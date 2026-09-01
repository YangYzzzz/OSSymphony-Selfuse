"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_image_dimensions__6fac91d7fdc9ac567c2522f5862ea0c3', 'check_gif_created__53c5a4efb8d4cc457e973819b6da6e9c', 'check_image_dimensions__f8142f5206c8e9d081140bc00d3fc065', 'check_image_grayscale__3e8ba2ce3f71bc35828deb0d6a741a53', 'check_gimp_theme_and_icons__48364d2e461d23c72442c1de0ff1aa61', 'check_image_dimensions__73568344a1f5a1b5f117e9d5ad98afd8', 'check_image_dimensions__7532d635f451c203dc4c9c0128dc01d3', 'check_image_dimensions__7d65af646c8700a29016dc9b48791bd4', 'check_image_mode__a066a98104bf8bbb65a3648dc0f66596', 'check_jpglist_content__0701a99d9d588eb9055008397d86fa70', 'check_image_color_mode__742fdbd7c8ad904b152a719424c640d8', 'check_image_file__9f8821d93c1286af4ab340177f286f49', 'check_image_files_count__c59a92777846cb030643927ac9628267', 'check_image_creation__3254fa401a56ff7bd854b02b7cc3c8d5', 'check_all_images_copied__8edb4c1b4d404d9ad51658001991b9be', 'check_image_dimensions__4cd9609153e6f966f651b58aa2377fc7', 'check_jpeg_export__eecb20058330e9d3e08bcd27beed34f8', 'check_gimp_export__c940821317f4393b9e9b23183cb5bab0', 'check_image_moved__ebf21364e626ee60bab03f3d9d12deff', 'check_file_exists_png__fda9115554280abf401514b41e43c7b7', 'check_gimp_multi_config__06220b8da042a0d79908bcbdd2480275', 'check_image_resized__d195f7f7ed997b32944186b135a8c2d3', 'check_gimp_multi_kv__814a465f74b3e354cb4f73bbafb3b9c5', 'check_image_dimensions__d91421f1c316cdb714dbaed12454e46f', 'check_image_removed__07113500ab2edd0ce92ed5a5e9e8b1aa', 'check_image_dimensions__ebe9791a027ef956f5ed9c5e287e6589', 'check_image_dimensions__96e0bf39947996b933028cd995042bff', 'check_image_dimensions__58760a5b814fae5bb8129a050acf3bc8', 'check_image_bottom_center__b7ad7d38cb36c439bba7c5f1f082cb99', 'check_gimp_png_dimensions__a3a90a0c275eafc9a08c6844b32106ee_qw35sft2_e2b530ec', 'check_gimp_undo_and_theme__402f86b3aa5767626ba0a47b48426a4d_qw35sft2_318d451c', 'check_gimp_image_scaled__6daba05d23c98408363b21827c57395a_qw35sft2_4ac4c429', 'check_gimp_grayscale_export__e4220acd868ea9e3c374fad54e986f17_qw35sft2_f93e3161', 'check_gimp_multi_config__4dfa9dfa61da7c0056a83284d2f6864c_qw35sft2_d847689f', 'check_image_mode__9561abdc46f323aa49f39811909461f6_qw35sft2_0e9f6226', 'check_gimp_img_dims__ee2a931afc08caeb2f7cbb79edbe900c_qw35sft2_fe1779ac', 'check_gimp_layer_names__f6dc05e0a2477162d5d2aaf28dcb399e_qw35sft2_f08733bf', 'check_gimp_mode_and_size__fc1506ee66b785bd74f37812296aa386_qw35sft2_da3275e4', 'check_image_size_exact__e32eea47feb3c826ca65513ea1a8c80e_qw35sft2_e9c9c771', 'check_gimp_darktable_plugin__93f45a58587e292e529b56db6ae27226_qw35sft2_57943c1b', 'check_gimp_saturation_and_contrast__918f3cc47704c578c4667b1ca6019596_qw35sft2_bf56261c', 'check_gimp_config_multi__d66c76f0f1e5e42f832534b4f9451e5c_qw35sft2_d46e6308', 'check_gimp_rotate_and_brightness__0ffd5d6c48e00cc11530a6062d573f43_qw35sft2_88dff48c', 'check_gimp_grayscale__a22a6a74169a7c2b7bd4f15c461d97a5_qw35sft2_8f0879ef', 'check_gimp_hflip__4ddc2b8f762091345213742b2e539174_qw35sft2_2d7d11a7', 'check_gimp_session_single_window__940ba8eb8583aa4e66192a1528358acd_qw35sft2_5a28ffa6', 'check_gimp_grayscale__b63a2518d07e9d19e5e5bc7d70fabd81_qw35sft2_490913b1', 'check_triangle_right_half__466f916fdda48299f33abc0cff1e195a_qw35sft2_8561ca43', 'check_gimp_image_size__ffabe808aeb093a1184f15812c26092c_qw35sft2_c2123506', 'check_gimp_image_and_layer__5e787196d91537456fde84b3a66eeecb_qw35sft2_604be911', 'check_gimp_bg_color__4d0454ba1a7fe5978c1747f21d797308_qw35sft2_3a1bf59f', 'check_gimp_hd_png__6d65b36cd0a438b39e8f21308e345a4a_qw35sft2_7958e00b', 'check_gimp_image_size__08c1ae085b68de96a31c7d77e1455141_qw35sft2_315c115c', 'check_gimp_layer_name__ad9b8a9ed8fc4caff86597f320aba5f2_qw35sft2_8b9cae8d', 'check_gimp_layer_name__f3ec988a99b0f642c7ac2a56df923954_qw35sft2_46923194', 'check_gimp_rot180__359fadf8cec20093badb6005d3119d0c_qw35sft2_789cc460', 'check_gimp_multi_config__4e9d034010f9da35cee4aabac6b48af6_qw35sft2_e8025bcd', 'check_gimp_running__be67663f10a797bb181cd6fa722efa34_qw35sft2_73fe97d0', 'check_gimp_saturation_and_size__e189188d82a47045e4921db4c0cbb19b_qw35sft2_df7145bb', 'check_gimp_grayscale_and_brightness__72b5517a5a800033bb98cd0fbc3da308_qw35sft2_96e118e2', 'check_gimp_mode_and_size__e5dd5bd57fb559e638682724445a4fb8_qw35sft2_4435fbc7', 'check_jpeg_and_contrast__99df0d1fa420bfe8f4d58bfec01f5388_qw35sft2_c9222dcb', 'check_image_dimensions__1649a3c87969acc3b9b0ad261438802f_qw35sft2_c01acf2c', 'check_gimp_dimensions__2254c4ac73b2e0e2d0bc95e32c93006c_qw35sft2_6addc5a5', 'check_gimp_image_grayscale__c8310c6f68cacbb7f806947f886f9156_qw35sft2_7b90abe5', 'check_gimp_grayscale__5cf8f74046cc654cd234f12e0bd56fa2_qw35sft2_7933c75d', 'check_gimp_icon_theme__8b7e0a81347ead101665acba492eaaa3_qw35sft2_a52da34f', 'check_gimp_image_rotated__65c7fe1c7d5d90a17513cb0cbcf64984_qw35sft2_dd6ec0d1', 'check_gimp_layer_bbox__6fac91d7fdc9ac567c2522f5862ea0c3_qw35sft2_a0ba7ccd', 'gimp_png_exported__2034cd2081886ce4d15c41eaddab9687_qw35sft2_efa8a2ed', 'check_gimp_bg_color__c69df4a705c1504ceebda9205ba51ae8_qw35sft2_16184163', 'check_gimp_hflip__e47abffb5b090802d85e5308a4df8da8_qw35sft2_16db1c14', 'check_gimp_theme__a2417c567c82dd8e5e291c848fcf86af_qw35sft2_b3a42674', 'check_gimp_saturation_increase__6b90b70b3d1b161e089c6f4f8582392c_qw35sft2_af0f3398', 'check_gimp_canvas_and_layer__0d4c5e593127b30f4279483bc536e453_qw35sft2_c95df0aa', 'check_desktop_jpeg_exists__e78c01185325a38b78f35c525682fa2c_qw35sft2_5e9ed279', 'check_gimp_mode_and_colors__4728404e320da0722e3f0bbf27779384_qw35sft2_0abecff9', 'check_gimp_jpeg_and_brightness__3494183243a5e3cf64667682e9d5cf69_qw35sft2_7b27cdcf', 'check_xcf_on_desktop__30619c16638f3c14c2b868710d3b511a_qw35sft2_93947704', 'check_gimp_config_multi__59d11816792f7732bb8631499d5a3dca_qw35sft2_1b01115c', 'check_gimp_flipped__f5b72d9abaa33261d0da153e643d90eb_qw35sft2_3eeb1203', 'check_gimp_flipped__917488a366fb0a5d7f8700214812c742_qw35sft2_2dd42002', 'check_gimp_rotated_90cw__480c196a41fcbd8c06328eee76532eca_qw35sft2_1da68272', 'check_gimp_file_exists__9dede0efb5d54dc2e4df137a06f2c4ba_qw35sft2_bc274fc9', 'check_gimp_dark_theme__48287b5138535fb1248d29dc8ceabce5_qw35sft2_3b068fa5', 'check_image_dimensions__9d59aa3517eac66c5b7e40985af8d11e_qw35sft2_f4b3c6df', 'check_gimp_text_on_left__f282f693e16d842a8fd79d730244ff72_qw35sft2_91885b50', 'check_triangle_lower_right__569dedcfdc3e7d291a20f0e6aca641c0_qw35sft2_2eeccac8', 'check_image_width__2bd82e314cabe351963c3af4824b9b87_qw35sft2_fb488a24', 'check_gimp_layer_bbox_opacity__c6fafc0c96511b4f7b94b1c390c2f8bd_qw35sft2_0a45d62f', 'check_gimp_bg_color__3cdf6e4d2fe41e5dc87ba95dd55c83c0_qw35sft2_0fae1857', 'check_gimp_vflip__e3f424dc6a1401f7936bf6fab235a0b4_qw35sft2_12bbae4a', 'check_gimp_icon_theme__a072fad36e1102073ae062e8b644b246_qw35sft2_1e9fa508', 'check_gimp_mode_and_size__94a488f0af5ba53f580224747ebe24fe_qw35sft2_d0cd4dc9', 'check_gimp_multi_layers__5ffee09eda4febddd3875d5ca422ec9b_qw35sft2_a13f473e', 'check_gimp_dims_and_brightness__1abf0259addada4e9320a412de93a720_qw35sft2_203dc4f9', 'check_gimp_config_multi__d46e54f2b8fa1eafde00e6af910958bf_qw35sft2_8379da51', 'check_gimp_theme_dark__ff19a7440acd5edb92781b332298a214_qw35sft2_683ee47b', 'check_gimp_png_exported__fcf5b1c8c41f94f2ef23ac5bc8b1ca54_qw35sft2_889db4e6', 'check_gimp_saturation_and_sharpness__cc63708ed8a0e82770a04258c8e27328_qw35sft2_d2f1e622', 'check_gimp_docks_window__e872a1deec6ced899c5abf40a9c68092_qw35sft2_1aaa3746', 'check_gimp_image_square__ba207a15040c6bc16b5ad1659309ab17_qw35sft2_d7307369', 'check_triangle_lower_half__bc6d20fd4fedacbbf14f4629287f2cc1_qw35sft2_e6b7454d', 'check_gimp_file_size_under__5c243e9c006808209adf482ad7ac1fe7_qw35sft2_af4a00b2', 'check_jpeg_file_exists__132eea53d5785c870b331091d6b9fa18_qw35sft2_89f22164', 'check_gimp_image_width__392ddc5ce54a4061c2326df12808cf9c_qw35sft2_12953686', 'check_gimp_layer_bbox__4cd9609153e6f966f651b58aa2377fc7_qw35sft2_ec87ac56', 'check_gimp_bg_color__06570640015ea4861f9ec59a895e97a7_qw35sft2_8a05249d', 'check_jpeg_exported__d6fa459b37f87da875d5b6370f3cc58a_qw35sft2_ba9ba93d', 'check_gimp_image_mode__f54b84b51bb414e1baf7bdf3d2e0e400_qw35sft2_5ce6c26e', 'check_gimp_layer_name__413739fbbe273da47f148920d82ae49b_qw35sft2_1ad7f5b9', 'check_gimp_mode_and_size__857686b0194013b75de5a057cd91e5ad_qw35sft2_7b9d3fdb', 'check_image_list_created__1dbc4bff650a401c3959c27fb5e70015_qw35sft2_3181279e', 'check_transition_and_png__dd35705ee09fd488827a6afdd7cbab81_qw35sft2_86c528ba', 'check_image_bottom_right__f123db274ed42f1e244c71177f056c09_qw35sft2_13a9b4f7', 'check_image_on_right__2e5aa88aaf11fa4bda2a55acfcab9dc5_qw35sft2_52f1b487', 'check_image_and_alignment__b349befcb19ab6c9d751d2833a5ceb8c_qw35sft2_e63d79de', 'check_image_and_pdf__cba3b2d7a3f4e1f51b9a84ac71e8e772_qw35sft2_7bb009ae', 'check_image_dimensions__b392913b04d99d3e7513b78f2dedf15d_qw35sft2_f9bcf566', 'check_image_compress_resize__8cc8ca1163a74ceeec7ab0fdb85713f3_qw35sft2_192625bb', 'check_image_resized_half__89f1952738bca9658039f047c420ba2d_qw35sft2_5c4c32f4', 'check_jpeg_compressed__19b267768d7fcf0b434b3cca9c02b15d_qw35sft2_d7a03f53', 'check_jpg_move__7427978e92f6fc0e5652a4713261a5a8_qw35sft2_fe0d3627', 'check_jpg_png_copy__ec3ddc36152b3d2fb4aee4f74f969e31_qw35sft2_14d06f7e', 'check_jpg_copy_with_count__58fb65f54a1152c1f3aecb552e15d932_qw35sft2_4feb954d', 'check_jpg_filelist__b885c2e3b122c9010091a38454018836_qw35sft2_eb1516da', 'check_vacation_jpg_copy__91c6f86e45abe96db716a4e3d1072be2_qw35sft2_07c5ddf5', 'check_vlc_image_adjust__33b1c5d9e144110cf196db624efc6d81_qw35sft2_33445adf']

def check_image_dimensions__6fac91d7fdc9ac567c2522f5862ea0c3(result, expected, **options):
    """Check image dimensions with partial credit for width and height."""
    if isinstance(result, str) or not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'target_width' in expected:
        total_checks += 1
        actual = result.get('width')
        if actual is not None and abs(actual - expected['target_width']) <= 2:
            score += 1.0
    if 'target_height' in expected:
        total_checks += 1
        actual = result.get('height')
        if actual is not None and abs(actual - expected['target_height']) <= 2:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_gif_created__53c5a4efb8d4cc457e973819b6da6e9c(result, expected, **options):
    """Check if a valid GIF file was created from the video.
    Partial credit:
      0.5 - File exists at expected path
      1.0 - File exists and is a valid GIF with reasonable dimensions
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.5
    if result.get('is_gif', False):
        width = result.get('width', 0)
        height = result.get('height', 0)
        min_size = expected.get('min_file_size', 1000)
        file_size = result.get('file_size', 0)
        if file_size >= min_size and width > 0 and (height > 0):
            score = 1.0
    return score

def check_image_dimensions__f8142f5206c8e9d081140bc00d3fc065(result, expected, **options):
    """Check layer dimensions with partial credit for width and height."""
    if isinstance(result, str) or not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    total_checks = 0
    tolerance = expected.get('tolerance', 2)
    if 'target_width' in expected:
        total_checks += 1
        actual = result.get('width')
        if actual is not None and abs(actual - expected['target_width']) <= tolerance:
            score += 1.0
    if 'target_height' in expected:
        total_checks += 1
        actual = result.get('height')
        if actual is not None and abs(actual - expected['target_height']) <= tolerance:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_image_grayscale__3e8ba2ce3f71bc35828deb0d6a741a53(result, expected, **options):
    """Check if the image is grayscale. Returns 1.0 if grayscale, 0.0 otherwise."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if isinstance(result, dict) and result.get('is_grayscale', False):
        return 1.0
    return 0.0

def check_gimp_theme_and_icons__48364d2e461d23c72442c1de0ff1aa61(result, expected, **options):
    """Check both theme and icon-theme with partial credit.

    Scoring: 0.5 for correct theme, 0.5 for correct icon-theme.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_theme = expected.get('expected_theme')
    expected_icon_theme = expected.get('expected_icon_theme')
    actual_theme = result.get('theme') if isinstance(result, dict) else None
    actual_icon_theme = result.get('icon_theme') if isinstance(result, dict) else None
    if expected_theme and actual_theme:
        if actual_theme.lower() == expected_theme.lower():
            score += 0.5
    if expected_icon_theme and actual_icon_theme:
        if actual_icon_theme.lower() == expected_icon_theme.lower():
            score += 0.5
    return min(score, 1.0)

def check_image_dimensions__73568344a1f5a1b5f117e9d5ad98afd8(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    expected_w = expected.get('expected_width')
    expected_h = expected.get('expected_height')
    if actual_w is None or actual_h is None:
        return 0.0
    score = 0.0
    if actual_w == expected_w:
        score += 0.5
    if actual_h == expected_h:
        score += 0.5
    return score

def check_image_dimensions__7532d635f451c203dc4c9c0128dc01d3(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    expected_w = expected.get('expected_width')
    expected_h = expected.get('expected_height')
    if actual_w is None or actual_h is None:
        return 0.0
    score = 0.0
    if actual_w == expected_w:
        score += 0.5
    if actual_h == expected_h:
        score += 0.5
    return score

def check_image_dimensions__7d65af646c8700a29016dc9b48791bd4(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    expected_w = expected.get('expected_width')
    expected_h = expected.get('expected_height')
    if actual_w is None or actual_h is None:
        return 0.0
    score = 0.0
    if actual_w == expected_w:
        score += 0.5
    if actual_h == expected_h:
        score += 0.5
    return score

def check_image_mode__a066a98104bf8bbb65a3648dc0f66596(result, expected, **options):
    """Check if image mode matches expected grayscale state."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_grayscale = expected.get('is_grayscale', True)
    actual_grayscale = result.get('is_grayscale', False)
    if expected_grayscale == actual_grayscale:
        return 1.0
    return 0.0

def check_jpglist_content__0701a99d9d588eb9055008397d86fa70(result, expected, **options):
    """Check if jpglist.txt contains correct sorted jpg filenames.

    Args:
        result: string output from cat command (file contents)
        expected: dict with 'expected_filenames' list, 'should_be_sorted', 'should_exclude_png'

    Returns:
        float: partial credit score 0.0-1.0
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    if isinstance(result, str):
        lines = [line.strip() for line in result.strip().split('\n') if line.strip()]
    else:
        return 0.0
    expected_filenames = expected.get('expected_filenames', [])
    should_be_sorted = expected.get('should_be_sorted', True)
    should_exclude_png = expected.get('should_exclude_png', True)
    score = 0.0
    found_count = 0
    for fname in expected_filenames:
        if any((fname in line for line in lines)):
            found_count += 1
    if expected_filenames:
        presence_score = found_count / len(expected_filenames)
        score += 0.5 * presence_score
    if should_be_sorted and lines:
        basenames = []
        for line in lines:
            parts = line.rsplit('/', 1)
            basenames.append(parts[-1] if parts else line)
        if basenames == sorted(basenames, key=str.lower):
            score += 0.25
    if should_exclude_png:
        has_png = any(('.png' in line.lower() for line in lines))
        if not has_png:
            score += 0.25
    return min(score, 1.0)

def check_image_color_mode__742fdbd7c8ad904b152a719424c640d8(result, expected, **options):
    """Check if the image color mode matches expected mode.

    Grayscale images in GIMP exported as PNG have mode 'L' (luminance).
    When exported as JPEG, grayscale images may still be 'RGB' but with
    all channels equal. We check for 'L' mode as primary, and also
    accept 'LA' (grayscale with alpha).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_mode = result.get('mode', '')
    expected_mode = expected.get('expected_mode', 'L')
    if actual_mode == expected_mode:
        return 1.0
    if expected_mode == 'L' and actual_mode in ('LA', 'P'):
        return 1.0
    return 0.0

def check_image_file__9f8821d93c1286af4ab340177f286f49(result, expected, **options):
    """Check if image file exists and has minimum size."""
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    score += 0.7
    min_size = expected.get('min_file_size', 100)
    if result.get('file_size', 0) >= min_size:
        score += 0.3
    return min(score, 1.0)

def check_image_files_count__c59a92777846cb030643927ac9628267(result, expected, **options):
    """Check that the expected number of image files exist in the directory."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    actual_count = result.get('count', 0)
    min_count = expected.get('min_count', 1)
    expected_count = expected.get('expected_count', 4)
    if actual_count == 0:
        return 0.0
    if actual_count >= expected_count:
        return 1.0
    if actual_count >= min_count:
        return actual_count / expected_count
    return 0.0

def check_image_creation__3254fa401a56ff7bd854b02b7cc3c8d5(result, expected, **options):
    """Check image exists with correct dimensions and color. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.3
    expected_w = expected.get('width')
    if expected_w and result.get('width') == expected_w:
        score += 0.25
    expected_h = expected.get('height')
    if expected_h and result.get('height') == expected_h:
        score += 0.25
    expected_color = expected.get('color', 'white')
    if expected_color == 'white' and result.get('is_white'):
        score += 0.2
    return min(score, 1.0)

def check_all_images_copied__8edb4c1b4d404d9ad51658001991b9be(result, expected, **options):
    """Check if all expected image files are present in the target directory.

    Args:
        result: dict from list_directory getter, representing directory tree
        expected: dict with 'expected_files' list of filenames

    Returns:
        float: partial credit score 0.0-1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_files = []
    if isinstance(result, dict):
        stack = [result]
        while stack:
            node = stack.pop()
            if isinstance(node, dict):
                for (key, value) in node.items():
                    if isinstance(value, dict):
                        stack.append(value)
                    else:
                        actual_files.append(key)
    score = 0.0
    per_file = 1.0 / len(expected_files)
    for f in expected_files:
        if f in actual_files:
            score += per_file
    return min(score, 1.0)

def check_image_dimensions__4cd9609153e6f966f651b58aa2377fc7(result, expected, **options):
    """Check layer dimensions with tolerance and partial credit.

    Checks target_height and target_width independently, awarding
    partial credit (0.5 each) so that a correct height but wrong width
    still gets 50%.
    """
    if isinstance(result, str) or not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    total_checks = 0
    tolerance = expected.get('tolerance', 2)
    if 'target_height' in expected:
        total_checks += 1
        actual_h = result.get('height')
        if actual_h is not None and abs(actual_h - expected['target_height']) <= tolerance:
            score += 1.0
    if 'target_width' in expected:
        total_checks += 1
        actual_w = result.get('width')
        if actual_w is not None and abs(actual_w - expected['target_width']) <= tolerance:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_jpeg_export__eecb20058330e9d3e08bcd27beed34f8(result, expected, **options):
    """Check if a valid JPEG file was exported. Partial credit for file existence."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.5
    if result.get('is_jpeg'):
        score += 0.5
    return score

def check_gimp_export__c940821317f4393b9e9b23183cb5bab0(result, expected, **options):
    """
    Check GIMP image export with partial credit:
    - 0.5 for file existing at the expected path
    - 0.5 for correct image dimensions (width and height)
    Returns float 0.0-1.0.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is not None and actual_height is not None:
        width_match = actual_width == expected_width if expected_width is not None else True
        height_match = actual_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            score += 0.5
    return min(score, 1.0)

def check_image_moved__ebf21364e626ee60bab03f3d9d12deff(result, expected, **options):
    """Check if the largest image on slide 1 has been moved to expected position."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    images = result.get('images', [])
    if not images:
        return 0.0
    largest = images[0]
    expected_x = expected.get('expected_left_cm', 0.0)
    expected_y = expected.get('expected_top_cm', 0.0)
    tolerance = expected.get('tolerance_cm', 0.3)
    score = 0.0
    actual_x = largest.get('left_cm', float('inf'))
    actual_y = largest.get('top_cm', float('inf'))
    if abs(actual_x - expected_x) <= tolerance:
        score += 0.5
    if abs(actual_y - expected_y) <= tolerance:
        score += 0.5
    return score

def check_file_exists_png__fda9115554280abf401514b41e43c7b7(result, expected, **options):
    """Check if file exists and is a valid PNG image."""
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.5
    first_bytes = result.get('first_bytes', '')
    if first_bytes.startswith('89504e47'):
        score += 0.3
    if result.get('size', 0) > 1024:
        score += 0.2
    return min(score, 1.0)

def check_gimp_multi_config__06220b8da042a0d79908bcbdd2480275(result, expected, **options):
    """Check multiple GIMP config values with partial credit.

    Args:
        result: dict with gimprc key-value pairs from getter
        expected: dict with 'checks' list of {key, value} dicts

    Returns:
        float: 0.0-1.0 score with partial credit per check
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    score_per_check = 1.0 / len(checks)
    total_score = 0.0
    for check in checks:
        key = check.get('key', '')
        expected_value = str(check.get('value', ''))
        actual_value = result.get(key)
        if actual_value is not None and str(actual_value) == expected_value:
            total_score += score_per_check
    return min(total_score, 1.0)

def check_image_resized__d195f7f7ed997b32944186b135a8c2d3(result, expected, **options):
    """Check if any image on slide 1 has been resized to the expected width."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    images = result.get('images', [])
    if not images:
        return 0.0
    expected_width = expected.get('expected_width_cm', 20.0)
    tolerance = expected.get('tolerance_cm', 0.5)
    for img in images:
        w = img.get('width_cm', 0)
        if abs(w - expected_width) <= tolerance:
            return 1.0
    return 0.0

def check_gimp_multi_kv__814a465f74b3e354cb4f73bbafb3b9c5(actual_config_path, rule):
    """Check multiple key-value pairs in GIMP config file with partial credit.

    Each check in rule['checks'] is verified against the config file.
    Score is evenly distributed across all checks.
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
    except Exception:
        return 0.0
    checks = rule.get('checks', [])
    if not checks:
        return 0.0
    score = 0.0
    weight = 1.0 / len(checks)
    for check in checks:
        key = check.get('key', '')
        value = check.get('value', '')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            items = line.strip().lstrip('(').rstrip(')\n').split()
            if len(items) >= 2 and items[0] == key and (items[-1] == value):
                score += weight
                break
    return min(score, 1.0)

def check_image_dimensions__d91421f1c316cdb714dbaed12454e46f(result, expected, **options):
    """Check if image dimensions match expected width and height. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('width') == expected.get('expected_width'):
        score += 0.5
    if result.get('height') == expected.get('expected_height'):
        score += 0.5
    return score

def check_image_removed__07113500ab2edd0ce92ed5a5e9e8b1aa(result, expected, **options):
    """Check if the large background image was removed from the slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_large = expected.get('expected_large_images', 0)
    actual_large = result.get('large_images', -1)
    if actual_large == expected_large:
        return 1.0
    original_large = expected.get('original_large_images', 1)
    if actual_large < original_large:
        return 0.5
    return 0.0

def check_image_dimensions__ebe9791a027ef956f5ed9c5e287e6589(result, expected, **options):
    """Check if image dimensions match expected width and height. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('width') == expected.get('expected_width'):
        score += 0.5
    if result.get('height') == expected.get('expected_height'):
        score += 0.5
    return score

def check_image_dimensions__96e0bf39947996b933028cd995042bff(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        return 0.0
    score = 0.0
    if actual_width == expected_width:
        score += 0.5
    if actual_height == expected_height:
        score += 0.5
    return score

def check_image_dimensions__58760a5b814fae5bb8129a050acf3bc8(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        return 0.0
    score = 0.0
    if actual_width == expected_width:
        score += 0.5
    if actual_height == expected_height:
        score += 0.5
    return score

def check_image_bottom_center__b7ad7d38cb36c439bba7c5f1f082cb99(result, expected, **options):
    """Check if the image is positioned in the bottom-center area of the slide.

    Bottom-center means:
    - Horizontally centered (center_x within 15% of slide center)
    - Vertically in bottom half (center_y in bottom 50% of slide)
    Partial credit: 0.5 for each condition met.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    slide_width = result.get('slide_width', 0)
    slide_height = result.get('slide_height', 0)
    center_x = result.get('center_x', 0)
    center_y = result.get('center_y', 0)
    if slide_width == 0 or slide_height == 0:
        return 0.0
    score = 0.0
    tolerance = expected.get('tolerance', 0.15)
    slide_center_x = slide_width / 2
    x_offset_ratio = abs(center_x - slide_center_x) / slide_width
    if x_offset_ratio <= tolerance:
        score += 0.5
    if center_y > slide_height * 0.5:
        score += 0.5
    return score

def check_gimp_png_dimensions__a3a90a0c275eafc9a08c6844b32106ee_qw35sft2_e2b530ec(result, expected, **options):
    """Check if exported PNG has the expected width and height. Partial credit 0.5 per dimension."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    exp_w = expected.get('width')
    exp_h = expected.get('height')
    score = 0.0
    if exp_w is not None and result.get('width') == exp_w:
        score += 0.5
    if exp_h is not None and result.get('height') == exp_h:
        score += 0.5
    return score

def check_gimp_undo_and_theme__402f86b3aa5767626ba0a47b48426a4d_qw35sft2_318d451c(result, expected, **options):
    """Check undo-levels and theme in GIMP config with partial credit (0.5 each)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_undo_levels = str(expected.get('undo_levels', '100')).strip()
    expected_theme = str(expected.get('theme', 'Dark')).strip().lower()
    actual_undo_levels = str(result.get('undo_levels') or '').strip()
    if actual_undo_levels == expected_undo_levels:
        score += 0.5
    actual_theme = str(result.get('theme') or '').strip().lower()
    if actual_theme == expected_theme:
        score += 0.5
    return score

def check_gimp_image_scaled__6daba05d23c98408363b21827c57395a_qw35sft2_4ac4c429(result, expected, **options):
    """Return 1.0 if the exported PNG matches the expected scaled dimensions (within tolerance)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_w = expected.get('expected_width')
    expected_h = expected.get('expected_height')
    tolerance = expected.get('tolerance', 20)
    actual_w = result.get('width', 0)
    actual_h = result.get('height', 0)
    if expected_w is None or expected_h is None:
        return 0.0
    w_ok = abs(actual_w - expected_w) <= tolerance
    h_ok = abs(actual_h - expected_h) <= tolerance
    return 1.0 if w_ok and h_ok else 0.0

def check_gimp_grayscale_export__e4220acd868ea9e3c374fad54e986f17_qw35sft2_f93e3161(result, expected, **options):
    """Check that export.jpg was saved in grayscale mode (PIL mode 'L')."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_mode = result.get('mode')
    if actual_mode is None:
        return 0.0
    expected_mode = expected.get('expected_mode', 'L')
    return 1.0 if actual_mode == expected_mode else 0.0

def check_gimp_multi_config__4dfa9dfa61da7c0056a83284d2f6864c_qw35sft2_d847689f(result, expected, **options):
    """Check multiple GIMP config settings with partial credit.

    result: local file path returned by gimp_config_file getter (str)
    expected: dict with key 'checks' mapping config-key -> expected-value
              e.g. {"checks": {"theme": "Blue", "icon-theme": "Symbolic-Inverted"}}
    """
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
    except Exception:
        return 0.0
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    num_checks = len(checks)
    per_check = 1.0 / num_checks
    score = 0.0
    for key, value in checks.items():
        pattern = f'\\({re.escape(key)}\\s+"({re.escape(str(value))})"\\)'
        if re.search(pattern, content):
            score += per_check
    return min(round(score, 4), 1.0)

def check_image_mode__9561abdc46f323aa49f39811909461f6_qw35sft2_0e9f6226(result, expected, **options):
    """Return 1.0 if the image's PIL mode is one of the expected modes."""
    if result is None or result.get('error'):
        return 0.0
    actual_mode = result.get('mode')
    expected_modes = expected.get('expected_modes', [])
    if actual_mode in expected_modes:
        return 1.0
    return 0.0

def check_gimp_img_dims__ee2a931afc08caeb2f7cbb79edbe900c_qw35sft2_fe1779ac(result, expected, **options):
    """
    Check that the image has exactly the expected pixel dimensions.
    'result' is the dict returned by get_gimp_img_dims__ee2a931afc08caeb2f7cbb79edbe900c.
    'expected' (already unwrapped from 'rules') must contain 'width' and 'height'.
    Returns 1.0 if both dimensions match, else 0.0.
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    expected_w = expected.get('width')
    expected_h = expected.get('height')
    actual_w = result.get('width')
    actual_h = result.get('height')
    if expected_w is not None and expected_h is not None:
        return 1.0 if actual_w == expected_w and actual_h == expected_h else 0.0
    return 0.0

def check_gimp_layer_names__f6dc05e0a2477162d5d2aaf28dcb399e_qw35sft2_f08733bf(result, expected, **options):
    """Check that all expected layer names exist in the XCF file. Partial credit per layer."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    layer_names = result.get('layer_names', [])
    target_names = expected.get('layer_names', [])
    if not target_names:
        return 0.0
    found = sum((1 for name in target_names if name in layer_names))
    return found / len(target_names)

def check_gimp_mode_and_size__fc1506ee66b785bd74f37812296aa386_qw35sft2_da3275e4(result, expected, **options):
    """Check image mode and dimensions with 0.5/0.5 partial credit.

    expected keys (unwrapped rules dict):
      expected_mode  - 'P' for indexed/palette
      expected_width - integer pixel width
      expected_height - integer pixel height
    """
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_mode = expected.get('expected_mode', 'P')
    if result.get('mode') == expected_mode:
        score += 0.5
    ew = expected.get('expected_width')
    eh = expected.get('expected_height')
    if ew is not None and eh is not None:
        if result.get('width') == ew and result.get('height') == eh:
            score += 0.5
    return min(score, 1.0)

def check_image_size_exact__e32eea47feb3c826ca65513ea1a8c80e_qw35sft2_e9c9c771(result, expected, **options):
    """Check if image dimensions match the expected width and height."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or not result.get('file_found', False):
        return 0.0
    actual_width = result.get('width', 0)
    actual_height = result.get('height', 0)
    exp_width = expected.get('width', -1)
    exp_height = expected.get('height', -1)
    if actual_width == exp_width and actual_height == exp_height:
        return 1.0
    return 0.0

def check_gimp_darktable_plugin__93f45a58587e292e529b56db6ae27226_qw35sft2_57943c1b(result, expected, **options):
    """Check that darktable is installed and its GIMP plugin file is present."""
    if not isinstance(result, dict):
        return 0.0
    pkg_installed = result.get('pkg_installed', False)
    plugin_present = result.get('plugin_present', False)
    if plugin_present and pkg_installed:
        return 1.0
    if pkg_installed:
        return 0.5
    return 0.0

def check_gimp_saturation_and_contrast__918f3cc47704c578c4667b1ca6019596_qw35sft2_bf56261c(result, expected, **options):
    """Partial credit: 0.5 for saturation increase, 0.5 for contrast increase."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    sat = result.get('saturation_mean', 0.0)
    contrast = result.get('contrast_std', 0.0)
    min_sat = expected.get('min_saturation', 0.42)
    min_contrast = expected.get('min_contrast_std', 70.0)
    if sat >= min_sat:
        score += 0.5
    if contrast >= min_contrast:
        score += 0.5
    return score

def check_gimp_config_multi__d66c76f0f1e5e42f832534b4f9451e5c_qw35sft2_d46e6308(result, expected, **options):
    """Check multiple GIMP config settings with partial credit (0.5 per check)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    per_check = 1.0 / len(checks)
    for key, expected_value in checks.items():
        actual_value = result.get(key, '')
        if actual_value == expected_value:
            score += per_check
    return min(score, 1.0)

def check_gimp_rotate_and_brightness__0ffd5d6c48e00cc11530a6062d573f43_qw35sft2_88dff48c(result, expected, **options):
    """
    Partial credit (0.5 each):
      - rotate_check: image dimensions are swapped from original 1099x730 to 730x1099
      - brightness_check: mean_brightness < brightness_threshold (original RGB mean ~71.74)
    Original image is 1099 wide x 730 tall; after 90-degree rotation: 730 wide x 1099 tall.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    width = result.get('width')
    height = result.get('height')
    mean_brightness = result.get('mean_brightness')
    expected_width = expected.get('expected_width', 730)
    expected_height = expected.get('expected_height', 1099)
    tolerance = expected.get('dim_tolerance', 2)
    brightness_threshold = expected.get('brightness_threshold', 65.0)
    if width is not None and height is not None and (abs(width - expected_width) <= tolerance) and (abs(height - expected_height) <= tolerance):
        score += 0.5
    if mean_brightness is not None and mean_brightness < brightness_threshold:
        score += 0.5
    return min(score, 1.0)

def check_gimp_grayscale__a22a6a74169a7c2b7bd4f15c461d97a5_qw35sft2_8f0879ef(result, expected, **options):
    """Check that the exported gate.jpeg is in Grayscale (mode 'L') mode."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_mode = result.get('mode')
    if actual_mode is None:
        return 0.0
    expected_mode = expected.get('expected_mode', 'L')
    return 1.0 if actual_mode == expected_mode else 0.0

def check_gimp_hflip__4ddc2b8f762091345213742b2e539174_qw35sft2_2d7d11a7(result, expected, **options):
    """Check that heron.jpeg was horizontally flipped by comparing thumbnail to flipped original."""
    from PIL import Image, ImageOps
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    original_path = expected.get('original_path')
    if not original_path:
        return 0.0
    actual_pixels = result.get('thumbnail_pixels')
    thumb_size = result.get('thumb_size', [64, 43])
    if not actual_pixels:
        return 0.0
    try:
        orig = Image.open(original_path).convert('RGB')
        flipped = ImageOps.mirror(orig)
        flipped_thumb = flipped.resize(tuple(thumb_size), Image.LANCZOS)
        expected_pixels = [list(p) for p in flipped_thumb.getdata()]
        if len(actual_pixels) != len(expected_pixels):
            return 0.0
        total_diff = sum((sum((abs(a - b) for a, b in zip(ap, ep))) for ap, ep in zip(actual_pixels, expected_pixels)))
        avg_diff = total_diff / (len(actual_pixels) * 3)
        return 1.0 if avg_diff < 20.0 else 0.0
    except Exception:
        return 0.0

def check_gimp_session_single_window__940ba8eb8583aa4e66192a1528358acd_qw35sft2_5a28ffa6(result, expected, **options):
    """Return 1.0 if GIMP single-window mode is disabled (value 'no' or key absent)."""
    if not isinstance(result, dict):
        return 0.0
    actual = result.get('single_window_mode', 'unknown')
    expected_val = expected.get('single_window_mode', 'no')
    if actual == expected_val:
        return 1.0
    if expected_val == 'no' and actual == 'unknown':
        return 1.0
    return 0.0

def check_gimp_grayscale__b63a2518d07e9d19e5e5bc7d70fabd81_qw35sft2_490913b1(result, expected, **options):
    """Check that the exported PNG is in Grayscale mode ('L' or 'LA')."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_mode = result.get('mode')
    if actual_mode is None:
        return 0.0
    expected_mode = expected.get('expected_mode', 'L')
    return 1.0 if actual_mode in (expected_mode, 'LA') else 0.0

def check_triangle_right_half__466f916fdda48299f33abc0cff1e195a_qw35sft2_8561ca43(result, expected, **options):
    """Check that the yellow triangle centroid is in the right half (centroid_x > canvas_width/2)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    cx = result.get('centroid_x')
    if cx is None:
        return 0.0
    canvas_width = result.get('width', 800)
    threshold = expected.get('threshold', canvas_width / 2)
    return 1.0 if cx > threshold else 0.0

def check_gimp_image_size__ffabe808aeb093a1184f15812c26092c_qw35sft2_c2123506(result, expected, **options):
    """Check if image dimensions match expected width and height."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    expected_w = expected.get('width', 0)
    expected_h = expected.get('height', 0)
    score = 0.0
    if expected_w and width == expected_w:
        score += 0.5
    if expected_h and height == expected_h:
        score += 0.5
    return score

def check_gimp_image_and_layer__5e787196d91537456fde84b3a66eeecb_qw35sft2_604be911(result, expected, **options):
    """
    Check two sub-goals with partial credit:
      0.5 - full image canvas dimensions match (verifies Scale Image operation)
      0.5 - non-transparent content bbox matches (verifies dog layer was pre-scaled)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tol = expected.get('tolerance', 2)
    exp_fw = expected.get('expected_file_w')
    exp_fh = expected.get('expected_file_h')
    actual_fw = result.get('file_w')
    actual_fh = result.get('file_h')
    if actual_fw is not None and actual_fh is not None:
        fw_ok = exp_fw is None or abs(actual_fw - exp_fw) <= tol
        fh_ok = exp_fh is None or abs(actual_fh - exp_fh) <= tol
        if fw_ok and fh_ok:
            score += 0.5
    exp_bw = expected.get('expected_bbox_w')
    exp_bh = expected.get('expected_bbox_h')
    actual_bw = result.get('bbox_w')
    actual_bh = result.get('bbox_h')
    if actual_bw is not None and actual_bh is not None:
        bw_ok = exp_bw is None or abs(actual_bw - exp_bw) <= tol
        bh_ok = exp_bh is None or abs(actual_bh - exp_bh) <= tol
        if bw_ok and bh_ok:
            score += 0.5
    return score

def check_gimp_bg_color__4d0454ba1a7fe5978c1747f21d797308_qw35sft2_3a1bf59f(result, expected, **options):
    """Check if background pixel color matches expected RGB values within tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 10)
    exp_r = expected.get('r', 0)
    exp_g = expected.get('g', 255)
    exp_b = expected.get('b', 0)
    actual_r = result.get('r', -1)
    actual_g = result.get('g', -1)
    actual_b = result.get('b', -1)
    if actual_r < 0 or actual_g < 0 or actual_b < 0:
        return 0.0
    if abs(actual_r - exp_r) <= tolerance and abs(actual_g - exp_g) <= tolerance and (abs(actual_b - exp_b) <= tolerance):
        return 1.0
    return 0.0

def check_gimp_hd_png__6d65b36cd0a438b39e8f21308e345a4a_qw35sft2_7958e00b(result, expected, **options):
    """Check if the exported PNG exists and matches expected width and height."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    exp_w = expected.get('width')
    exp_h = expected.get('height')
    score = 0.0
    if exp_w is not None and result.get('width') == exp_w:
        score += 0.5
    if exp_h is not None and result.get('height') == exp_h:
        score += 0.5
    return score

def check_gimp_image_size__08c1ae085b68de96a31c7d77e1455141_qw35sft2_315c115c(result, expected, **options):
    """Check that the exported image has the expected pixel dimensions."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if actual_width == expected_width and actual_height == expected_height:
        return 1.0
    return 0.0

def check_gimp_layer_name__ad9b8a9ed8fc4caff86597f320aba5f2_qw35sft2_8b9cae8d(result, expected, **options):
    """Check that a specific layer name exists in the XCF file."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    layer_names = result.get('layer_names', [])
    target = expected.get('layer_name', '')
    if not target:
        return 0.0
    return 1.0 if target in layer_names else 0.0

def check_gimp_layer_name__f3ec988a99b0f642c7ac2a56df923954_qw35sft2_46923194(result, expected, **options):
    """Stub: superseded by check_gimp_layer_name__ad9b8a9ed8fc4caff86597f320aba5f2."""
    return 0.0

def check_gimp_rot180__359fadf8cec20093badb6005d3119d0c_qw35sft2_789cc460(result, expected, **options):
    """
    Check that the exported image is a 180-degree rotation of the original.
    'result' is the dict returned by get_gimp_rot180_check__359fadf8cec20093badb6005d3119d0c.
    'expected' (already unwrapped from 'rules') must contain 'threshold'.
    Returns 1.0 if similarity >= threshold, else 0.0.
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    similarity = result.get('similarity', 0.0)
    threshold = expected.get('threshold', 0.97)
    return 1.0 if similarity >= threshold else 0.0

def check_gimp_multi_config__4e9d034010f9da35cee4aabac6b48af6_qw35sft2_e8025bcd(result, expected, **options):
    """Check multiple GIMP config settings with partial credit.

    result: local file path returned by gimp_config_file getter (str)
    expected: dict with key 'checks' mapping config-key -> expected-value
              e.g. {"checks": {"theme": "Blue", "icon-theme": "Legacy"}}
    """
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
    except Exception:
        return 0.0
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    num_checks = len(checks)
    per_check = 1.0 / num_checks
    score = 0.0
    for key, value in checks.items():
        pattern = f'\\({re.escape(key)}\\s+"({re.escape(str(value))})"\\)'
        if re.search(pattern, content):
            score += per_check
    return min(round(score, 4), 1.0)

def check_gimp_running__be67663f10a797bb181cd6fa722efa34_qw35sft2_73fe97d0(result, expected, **options):
    """Return 1.0 if GIMP is currently running, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('running', False) else 0.0

def check_gimp_saturation_and_size__e189188d82a47045e4921db4c0cbb19b_qw35sft2_df7145bb(result, expected, **options):
    """Partial credit: 0.5 for saturation increase, 0.5 for correct image dimensions."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    sat = result.get('saturation_mean', 0.0)
    width = result.get('width', 0)
    height = result.get('height', 0)
    min_sat = expected.get('min_saturation', 0.42)
    expected_width = expected.get('expected_width', 800)
    expected_height = expected.get('expected_height', 600)
    if sat >= min_sat:
        score += 0.5
    if width == expected_width and height == expected_height:
        score += 0.5
    return score

def check_gimp_grayscale_and_brightness__72b5517a5a800033bb98cd0fbc3da308_qw35sft2_96e118e2(result, expected, **options):
    """
    Partial credit (0.5 each):
      - grayscale_check: image mode indicates grayscale ('L' or 'LA')
      - brightness_check: mean grayscale brightness < brightness_threshold (original grayscale mean ~71.65)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    mode = result.get('mode')
    mean_brightness = result.get('mean_brightness')
    expected_mode = expected.get('expected_mode', 'L')
    brightness_threshold = expected.get('brightness_threshold', 65.0)
    if mode is not None and mode in ('L', 'LA'):
        score += 0.5
    if mean_brightness is not None and mean_brightness < brightness_threshold:
        score += 0.5
    return min(score, 1.0)

def check_gimp_mode_and_size__e5dd5bd57fb559e638682724445a4fb8_qw35sft2_4435fbc7(result, expected, **options):
    """Check image mode and dimensions with 0.5/0.5 partial credit.

    expected keys (unwrapped rules dict):
      expected_mode  - 'P' for indexed/palette
      expected_width - integer pixel width
      expected_height - integer pixel height
    """
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_mode = expected.get('expected_mode', 'P')
    if result.get('mode') == expected_mode:
        score += 0.5
    ew = expected.get('expected_width')
    eh = expected.get('expected_height')
    if ew is not None and eh is not None:
        if result.get('width') == ew and result.get('height') == eh:
            score += 0.5
    return min(score, 1.0)

def check_jpeg_and_contrast__99df0d1fa420bfe8f4d58bfec01f5388_qw35sft2_c9222dcb(result, expected, **options):
    """Check that file is JPEG format AND contrast was increased vs original.

    Scoring: 0.5 for JPEG format, 0.5 for contrast increase.
    Original berries.png avg_std ~51.1; threshold 56.0 ensures meaningful contrast boost.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    img_format = result.get('format', '')
    avg_std = result.get('avg_std', 0.0)
    exp_format = expected.get('expected_format', 'JPEG')
    min_std = expected.get('min_contrast_std', 56.0)
    if img_format == exp_format:
        score += 0.5
    if avg_std >= min_std:
        score += 0.5
    return score

def check_image_dimensions__1649a3c87969acc3b9b0ad261438802f_qw35sft2_c01acf2c(result, expected, **options):
    """Return 1.0 if image width and height match the expected values."""
    if result is None or result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    exp_w = expected.get('width')
    exp_h = expected.get('height')
    if actual_w == exp_w and actual_h == exp_h:
        return 1.0
    return 0.0

def check_gimp_dimensions__2254c4ac73b2e0e2d0bc95e32c93006c_qw35sft2_6addc5a5(result, expected, **options):
    """Check that the exported gate.jpeg matches the expected dimensions (±5px tolerance)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    if actual_w is None or actual_h is None:
        return 0.0
    expected_w = expected.get('expected_width')
    expected_h = expected.get('expected_height')
    tolerance = expected.get('tolerance', 5)
    score = 0.0
    if expected_w is not None and abs(actual_w - expected_w) <= tolerance:
        score += 0.5
    if expected_h is not None and abs(actual_h - expected_h) <= tolerance:
        score += 0.5
    return score

def check_gimp_image_grayscale__c8310c6f68cacbb7f806947f886f9156_qw35sft2_7b90abe5(result, expected, **options):
    """Check if exported image is in grayscale mode (L or LA)."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') and result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    return 1.0 if mode in ('L', 'LA') else 0.0

def check_gimp_grayscale__5cf8f74046cc654cd234f12e0bd56fa2_qw35sft2_7933c75d(result, expected, **options):
    """Check that heron.jpeg was converted to grayscale (mode 'L')."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_mode = result.get('mode')
    if actual_mode is None:
        return 0.0
    expected_mode = expected.get('expected_mode', 'L')
    return 1.0 if actual_mode == expected_mode else 0.0

def check_gimp_icon_theme__8b7e0a81347ead101665acba492eaaa3_qw35sft2_a52da34f(result, expected, **options):
    """Return 1.0 if GIMP icon theme is set to Symbolic."""
    if not isinstance(result, dict):
        return 0.0
    actual = result.get('icon_theme', 'unknown').lower().strip('"\'')
    expected_val = expected.get('icon_theme', 'symbolic').lower().strip('"\'')
    return 1.0 if expected_val in actual or actual == expected_val else 0.0

def check_gimp_image_rotated__65c7fe1c7d5d90a17513cb0cbcf64984_qw35sft2_dd6ec0d1(result, expected, **options):
    """Return 1.0 if the exported PNG is in portrait orientation (rotated 90 degrees)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('is_portrait') else 0.0

def check_gimp_layer_bbox__6fac91d7fdc9ac567c2522f5862ea0c3_qw35sft2_a0ba7ccd(result, expected, **options):
    """Check that exported PNG has the expected non-transparent content bounding box dimensions."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tol = expected.get('tolerance', 2)
    exp_w = expected.get('expected_bbox_w')
    exp_h = expected.get('expected_bbox_h')
    actual_w = result.get('bbox_w')
    actual_h = result.get('bbox_h')
    if actual_w is None or actual_h is None:
        return 0.0
    w_ok = exp_w is None or abs(actual_w - exp_w) <= tol
    h_ok = exp_h is None or abs(actual_h - exp_h) <= tol
    return 1.0 if w_ok and h_ok else 0.0

def gimp_png_exported__2034cd2081886ce4d15c41eaddab9687_qw35sft2_efa8a2ed(result, expected, **options):
    """Check that the exported PNG exists and has the expected dimensions."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    width = result.get('width')
    height = result.get('height')
    if width is None or height is None:
        return 0.0
    expected_w = expected.get('width', 800)
    expected_h = expected.get('height', 800)
    if width == expected_w and height == expected_h:
        return 1.0
    return 0.0

def check_gimp_bg_color__c69df4a705c1504ceebda9205ba51ae8_qw35sft2_16184163(result, expected, **options):
    """Check if background pixel color matches expected RGB values within tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 10)
    exp_r = expected.get('r', 0)
    exp_g = expected.get('g', 255)
    exp_b = expected.get('b', 0)
    actual_r = result.get('r', -1)
    actual_g = result.get('g', -1)
    actual_b = result.get('b', -1)
    if actual_r < 0 or actual_g < 0 or actual_b < 0:
        return 0.0
    if abs(actual_r - exp_r) <= tolerance and abs(actual_g - exp_g) <= tolerance and (abs(actual_b - exp_b) <= tolerance):
        return 1.0
    return 0.0

def check_gimp_hflip__e47abffb5b090802d85e5308a4df8da8_qw35sft2_16db1c14(result, expected, **options):
    """
    Check that the exported image is a horizontal flip of the original.
    'result' is the dict returned by get_gimp_hflip_check__e47abffb5b090802d85e5308a4df8da8.
    'expected' (already unwrapped from 'rules') must contain 'threshold'.
    Returns 1.0 if similarity >= threshold, else 0.0.
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    similarity = result.get('similarity', 0.0)
    threshold = expected.get('threshold', 0.97)
    return 1.0 if similarity >= threshold else 0.0

def check_gimp_theme__a2417c567c82dd8e5e291c848fcf86af_qw35sft2_b3a42674(result, expected, **options):
    """Check if GIMP theme matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_theme = expected.get('expected_theme', 'dark').lower()
    actual_theme = (result.get('theme') or '').lower()
    return 1.0 if actual_theme == expected_theme else 0.0

def check_gimp_saturation_increase__6b90b70b3d1b161e089c6f4f8582392c_qw35sft2_af0f3398(result, expected, **options):
    """Check that mean HSV saturation of the saved image exceeds the expected minimum."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    sat = result.get('saturation_mean', 0.0)
    min_sat = expected.get('min_saturation', 0.42)
    return 1.0 if sat >= min_sat else 0.0

def check_gimp_canvas_and_layer__0d4c5e593127b30f4279483bc536e453_qw35sft2_c95df0aa(result, expected, **options):
    """Partial credit: 0.5 for canvas.png on Desktop, 0.5 for 'Square' layer in XCF."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('png_exists'):
        score += 0.5
    if result.get('square_layer_found'):
        score += 0.5
    return score

def check_desktop_jpeg_exists__e78c01185325a38b78f35c525682fa2c_qw35sft2_5e9ed279(result, expected, **options):
    """Check that yicun.jpg was created on the Desktop as a valid JPEG."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('valid_jpeg'):
        return 1.0
    if result.get('exists'):
        return 0.5
    return 0.0

def check_gimp_mode_and_colors__4728404e320da0722e3f0bbf27779384_qw35sft2_0abecff9(result, expected, **options):
    """Check that the image is indexed mode and uses <= max_colors unique colors.

    Partial credit: 0.5 for correct mode, 0.5 for color count within limit.

    expected keys (unwrapped rules dict):
      expected_mode - 'P' (default, indexed/palette mode)
      max_colors    - integer maximum allowed unique colors (default 128)
    """
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_mode = expected.get('expected_mode', 'P')
    if result.get('mode') == expected_mode:
        score += 0.5
    max_colors = expected.get('max_colors', 128)
    color_count = result.get('color_count', -1)
    if color_count > 0 and color_count <= max_colors:
        score += 0.5
    return min(score, 1.0)

def check_gimp_jpeg_and_brightness__3494183243a5e3cf64667682e9d5cf69_qw35sft2_7b27cdcf(result, expected, **options):
    """
    Full credit (1.0):
      - dark_photo.jpg exists on Desktop with mean RGB brightness < brightness_threshold
    Threshold defaults to 68.0 (~5% reduction from original ~71.74), accepting any genuine
    brightness reduction. PNG overwrite check removed: instruction only requires JPEG export.
    """
    if not isinstance(result, dict):
        return 0.0
    brightness_threshold = expected.get('brightness_threshold', 68.0)
    jpeg_exists = result.get('jpeg_exists', False)
    jpeg_mean = result.get('jpeg_mean_brightness')
    if jpeg_exists and jpeg_mean is not None and (jpeg_mean < brightness_threshold):
        return 1.0
    return 0.0

def check_xcf_on_desktop__30619c16638f3c14c2b868710d3b511a_qw35sft2_93947704(result, expected, **options):
    """Return 1.0 if the expected XCF filename is found on Desktop with correct canvas dimensions."""
    if not isinstance(result, dict):
        return 0.0
    files = result.get('files', [])
    filename = expected.get('filename', 'new_image.xcf')
    if filename not in files:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is not None and result.get('width') != expected_width:
        return 0.0
    if expected_height is not None and result.get('height') != expected_height:
        return 0.0
    return 1.0

def check_gimp_config_multi__59d11816792f7732bb8631499d5a3dca_qw35sft2_1b01115c(result, expected, **options):
    """Check multiple GIMP config settings with partial credit (0.5 per check)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    per_check = 1.0 / len(checks)
    for key, expected_value in checks.items():
        actual_value = result.get(key, '')
        if actual_value == expected_value:
            score += per_check
    return min(score, 1.0)

def check_gimp_flipped__f5b72d9abaa33261d0da153e643d90eb_qw35sft2_3eeb1203(result, expected, **options):
    """Check that the exported dog_flipped.png matches a programmatically-flipped version of the original."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or not result.get('exists', False):
        return 0.0
    if result.get('flip_match', False):
        return 1.0
    return 0.0

def check_gimp_flipped__917488a366fb0a5d7f8700214812c742_qw35sft2_2dd42002(result, expected, **options):
    """Check that gate.jpeg was flipped horizontally.

    In the original image, left_half_avg ≈ 109.72 (brighter, stone wall side).
    After a horizontal flip, the darker iron-gate side moves to the left,
    so left_half_avg drops to ≈ 77.94.  The threshold of 95 cleanly separates
    the two states.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    left_avg = result.get('left_half_avg')
    if left_avg is None:
        return 0.0
    threshold = expected.get('left_half_avg_max', 95.0)
    return 1.0 if left_avg < threshold else 0.0

def check_gimp_rotated_90cw__480c196a41fcbd8c06328eee76532eca_qw35sft2_1da68272(result, expected, **options):
    """Check that heron.jpeg was rotated 90 degrees clockwise (dimensions swapped)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        return 0.0
    expected_width = expected.get('expected_width', 853)
    expected_height = expected.get('expected_height', 1280)
    if actual_width == expected_width and actual_height == expected_height:
        return 1.0
    return 0.0

def check_gimp_file_exists__9dede0efb5d54dc2e4df137a06f2c4ba_qw35sft2_bc274fc9(result, expected, **options):
    """Check that export.jpg was created in the Documents folder."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('exists') and result.get('size', 0) > 0 else 0.0

def check_gimp_dark_theme__48287b5138535fb1248d29dc8ceabce5_qw35sft2_3b068fa5(result, expected, **options):
    """Return 1.0 if GIMP theme is set to Dark."""
    if not isinstance(result, dict):
        return 0.0
    actual = result.get('theme', 'unknown').lower().strip('"\'')
    expected_val = expected.get('theme', 'dark').lower().strip('"\'')
    return 1.0 if expected_val in actual or actual == expected_val else 0.0

def check_image_dimensions__9d59aa3517eac66c5b7e40985af8d11e_qw35sft2_f4b3c6df(result, expected, **options):
    """Return 1.0 if image width and height match the expected values."""
    if result is None or result.get('error'):
        return 0.0
    actual_w = result.get('width')
    actual_h = result.get('height')
    exp_w = expected.get('width')
    exp_h = expected.get('height')
    if actual_w == exp_w and actual_h == exp_h:
        return 1.0
    return 0.0

def check_gimp_text_on_left__f282f693e16d842a8fd79d730244ff72_qw35sft2_91885b50(result, expected, **options):
    """Return 1.0 if the text (dark pixels) in the PNG is on the left half of the canvas."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('is_on_left_half') else 0.0

def check_triangle_lower_right__569dedcfdc3e7d291a20f0e6aca641c0_qw35sft2_2eeccac8(result, expected, **options):
    """Check that the yellow triangle centroid is in the lower-right quadrant (X>500 AND Y>500)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    cx = result.get('centroid_x')
    cy = result.get('centroid_y')
    if cx is None or cy is None:
        return 0.0
    x_threshold = expected.get('x_threshold', 500)
    y_threshold = expected.get('y_threshold', 500)
    score = 0.0
    if cx > x_threshold:
        score += 0.5
    if cy > y_threshold:
        score += 0.5
    return score

def check_image_width__2bd82e314cabe351963c3af4824b9b87_qw35sft2_fb488a24(result, expected, **options):
    """Check that the exported image width matches the expected value."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_width = result.get('width')
    if actual_width is None:
        return 0.0
    expected_width = expected.get('expected_width', 400)
    tolerance = expected.get('tolerance', 2)
    return 1.0 if abs(actual_width - expected_width) <= tolerance else 0.0

def check_gimp_layer_bbox_opacity__c6fafc0c96511b4f7b94b1c390c2f8bd_qw35sft2_0a45d62f(result, expected, **options):
    """
    Check two sub-goals with partial credit:
      0.5 - non-transparent bbox matches expected dimensions
      0.5 - max alpha value indicates ~50% layer opacity
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tol = expected.get('tolerance', 2)
    exp_w = expected.get('expected_bbox_w')
    exp_h = expected.get('expected_bbox_h')
    actual_w = result.get('bbox_w')
    actual_h = result.get('bbox_h')
    if actual_w is not None and actual_h is not None:
        w_ok = exp_w is None or abs(actual_w - exp_w) <= tol
        h_ok = exp_h is None or abs(actual_h - exp_h) <= tol
        if w_ok and h_ok:
            score += 0.5
    exp_max_alpha = expected.get('expected_max_alpha', 128)
    alpha_tol = expected.get('alpha_tolerance', 30)
    actual_max_alpha = result.get('max_alpha')
    if actual_max_alpha is not None and abs(actual_max_alpha - exp_max_alpha) <= alpha_tol:
        score += 0.5
    return score

def check_gimp_bg_color__3cdf6e4d2fe41e5dc87ba95dd55c83c0_qw35sft2_0fae1857(result, expected, **options):
    """Check if background pixel color matches expected RGB values within tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 10)
    exp_r = expected.get('r', 255)
    exp_g = expected.get('g', 0)
    exp_b = expected.get('b', 0)
    actual_r = result.get('r', -1)
    actual_g = result.get('g', -1)
    actual_b = result.get('b', -1)
    if actual_r < 0 or actual_g < 0 or actual_b < 0:
        return 0.0
    if abs(actual_r - exp_r) <= tolerance and abs(actual_g - exp_g) <= tolerance and (abs(actual_b - exp_b) <= tolerance):
        return 1.0
    return 0.0

def check_gimp_vflip__e3f424dc6a1401f7936bf6fab235a0b4_qw35sft2_12bbae4a(result, expected, **options):
    """
    Check that the exported image is a vertical flip of the original.
    'result' is the dict returned by get_gimp_vflip_check__e3f424dc6a1401f7936bf6fab235a0b4.
    'expected' (already unwrapped from 'rules') must contain 'threshold'.
    Returns 1.0 if similarity >= threshold, else 0.0.
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    similarity = result.get('similarity', 0.0)
    threshold = expected.get('threshold', 0.97)
    return 1.0 if similarity >= threshold else 0.0

def check_gimp_icon_theme__a072fad36e1102073ae062e8b644b246_qw35sft2_1e9fa508(result, expected, **options):
    """Check if GIMP icon-theme matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_icon = expected.get('expected_icon_theme', 'symbolic').lower()
    actual_icon = (result.get('icon_theme') or '').lower()
    return 1.0 if actual_icon == expected_icon else 0.0

def check_gimp_mode_and_size__94a488f0af5ba53f580224747ebe24fe_qw35sft2_d0cd4dc9(result, expected, **options):
    """Check image mode and dimensions with 0.5/0.5 partial credit.

    expected keys (unwrapped rules dict):
      expected_mode  - 'L' for grayscale
      expected_width - integer pixel width
      expected_height - integer pixel height
    """
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_mode = expected.get('expected_mode', 'L')
    if result.get('mode') == expected_mode:
        score += 0.5
    ew = expected.get('expected_width')
    eh = expected.get('expected_height')
    if ew is not None and eh is not None:
        if result.get('width') == ew and result.get('height') == eh:
            score += 0.5
    return min(score, 1.0)

def check_gimp_multi_layers__5ffee09eda4febddd3875d5ca422ec9b_qw35sft2_a13f473e(result, expected, **options):
    """Check that two specific layer names both exist in the XCF. Partial credit 0.5 each."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    layer_names = result.get('layer_names', [])
    target_names = expected.get('layer_names', [])
    if not target_names:
        return 0.0
    found = sum((1 for name in target_names if name in layer_names))
    return found / len(target_names)

def check_gimp_dims_and_brightness__1abf0259addada4e9320a412de93a720_qw35sft2_203dc4f9(result, expected, **options):
    """
    Partial credit (0.5 each):
      - dims_check: image width == expected_width and height == expected_height (within tolerance)
      - brightness_check: mean_brightness < brightness_threshold (original RGB mean ~71.74)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    width = result.get('width')
    height = result.get('height')
    mean_brightness = result.get('mean_brightness')
    expected_width = expected.get('expected_width', 800)
    expected_height = expected.get('expected_height', 600)
    tolerance = expected.get('dim_tolerance', 2)
    brightness_threshold = expected.get('brightness_threshold', 65.0)
    if width is not None and height is not None and (abs(width - expected_width) <= tolerance) and (abs(height - expected_height) <= tolerance):
        score += 0.5
    if mean_brightness is not None and mean_brightness < brightness_threshold:
        score += 0.5
    return min(score, 1.0)

def check_gimp_config_multi__d46e54f2b8fa1eafde00e6af910958bf_qw35sft2_8379da51(result, expected, **options):
    """Check multiple GIMP config settings with partial credit (0.5 per check)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    per_check = 1.0 / len(checks)
    for key, expected_value in checks.items():
        actual_value = result.get(key, '')
        if actual_value == expected_value:
            score += per_check
    return min(score, 1.0)

def check_gimp_theme_dark__ff19a7440acd5edb92781b332298a214_qw35sft2_683ee47b(result, expected, **options):
    """Check if GIMP theme is set to Dark in gimprc config."""
    if not isinstance(result, dict):
        return 0.0
    config_content = result.get('config', '')
    if not config_content:
        return 0.0
    expected_theme = expected.get('theme', 'Dark')
    match = re.search('\\(theme\\s+"([^"]+)"\\)', config_content, re.IGNORECASE)
    if match:
        actual_theme = match.group(1)
        return 1.0 if actual_theme.lower() == expected_theme.lower() else 0.0
    return 0.0

def check_gimp_png_exported__fcf5b1c8c41f94f2ef23ac5bc8b1ca54_qw35sft2_889db4e6(result, expected, **options):
    """Check that a PNG file was exported to the Desktop."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    expected_exists = expected.get('expected_exists', True)
    actual_exists = result.get('has_png', False)
    return 1.0 if actual_exists == expected_exists else 0.0

def check_gimp_saturation_and_sharpness__cc63708ed8a0e82770a04258c8e27328_qw35sft2_d2f1e622(result, expected, **options):
    """Partial credit: 0.5 for saturation increase, 0.5 for sharpness increase."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    sat = result.get('saturation_mean', 0.0)
    laplacian_var = result.get('laplacian_var', 0.0)
    min_sat = expected.get('min_saturation', 0.42)
    min_laplacian = expected.get('min_laplacian_var', 110.0)
    if sat >= min_sat:
        score += 0.5
    if laplacian_var >= min_laplacian:
        score += 0.5
    return score

def check_gimp_docks_window__e872a1deec6ced899c5abf40a9c68092_qw35sft2_1aaa3746(result, expected, **options):
    """Partial credit: 0.5 for docks hidden, 0.5 for single-window mode disabled."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('hide_docks') == expected.get('hide_docks', 'yes'):
        score += 0.5
    actual_swm = result.get('single_window_mode', 'unknown')
    expected_swm = expected.get('single_window_mode', 'no')
    if actual_swm == expected_swm:
        score += 0.5
    elif expected_swm == 'no' and actual_swm == 'unknown':
        score += 0.5
    return score

def check_gimp_image_square__ba207a15040c6bc16b5ad1659309ab17_qw35sft2_d7307369(result, expected, **options):
    """Return 1.0 if the exported PNG is a square (equal width and height)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 0)
    w = result.get('width', -1)
    h = result.get('height', -1)
    if w <= 0 or h <= 0:
        return 0.0
    return 1.0 if abs(w - h) <= tolerance else 0.0

def check_triangle_lower_half__bc6d20fd4fedacbbf14f4629287f2cc1_qw35sft2_e6b7454d(result, expected, **options):
    """Check that the yellow triangle centroid is in the lower half (centroid_y > canvas_height/2)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    cy = result.get('centroid_y')
    if cy is None:
        return 0.0
    canvas_height = result.get('height', 800)
    threshold = expected.get('threshold', canvas_height / 2)
    return 1.0 if cy > threshold else 0.0

def check_gimp_file_size_under__5c243e9c006808209adf482ad7ac1fe7_qw35sft2_af4a00b2(result, expected, **options):
    """Check that the exported file is smaller than max_bytes (low quality export)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_size = result.get('size')
    if actual_size is None:
        return 0.0
    max_bytes = expected.get('max_bytes', 65536)
    return 1.0 if actual_size <= max_bytes else 0.0

def check_jpeg_file_exists__132eea53d5785c870b331091d6b9fa18_qw35sft2_89f22164(result, expected, **options):
    """Return 1.0 if a JPEG file was exported to the Desktop, else 0.0."""
    if result is None or (isinstance(result, dict) and 'error' in result):
        return 0.0
    jpeg_exists = result.get('jpeg_exists', False) if isinstance(result, dict) else False
    expected_val = expected.get('jpeg_exists', True)
    return 1.0 if bool(jpeg_exists) == bool(expected_val) else 0.0

def check_gimp_image_width__392ddc5ce54a4061c2326df12808cf9c_qw35sft2_12953686(result, expected, **options):
    """Check that heron.jpeg was scaled to the expected width."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_width = result.get('width')
    if actual_width is None:
        return 0.0
    expected_width = expected.get('expected_width', 640)
    return 1.0 if actual_width == expected_width else 0.0

def check_gimp_layer_bbox__4cd9609153e6f966f651b58aa2377fc7_qw35sft2_ec87ac56(result, expected, **options):
    """Check that exported PNG has the expected non-transparent content bounding box dimensions."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tol = expected.get('tolerance', 2)
    exp_w = expected.get('expected_bbox_w')
    exp_h = expected.get('expected_bbox_h')
    actual_w = result.get('bbox_w')
    actual_h = result.get('bbox_h')
    if actual_w is None or actual_h is None:
        return 0.0
    w_ok = exp_w is None or abs(actual_w - exp_w) <= tol
    h_ok = exp_h is None or abs(actual_h - exp_h) <= tol
    return 1.0 if w_ok and h_ok else 0.0

def check_gimp_bg_color__06570640015ea4861f9ec59a895e97a7_qw35sft2_8a05249d(result, expected, **options):
    """Check if background pixel color matches expected RGB values within tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 10)
    exp_r = expected.get('r', 0)
    exp_g = expected.get('g', 0)
    exp_b = expected.get('b', 255)
    actual_r = result.get('r', -1)
    actual_g = result.get('g', -1)
    actual_b = result.get('b', -1)
    if actual_r < 0 or actual_g < 0 or actual_b < 0:
        return 0.0
    if abs(actual_r - exp_r) <= tolerance and abs(actual_g - exp_g) <= tolerance and (abs(actual_b - exp_b) <= tolerance):
        return 1.0
    return 0.0

def check_jpeg_exported__d6fa459b37f87da875d5b6370f3cc58a_qw35sft2_ba9ba93d(result, expected, **options):
    """Check that the JPEG file was exported successfully."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('exists') else 0.0

def check_gimp_image_mode__f54b84b51bb414e1baf7bdf3d2e0e400_qw35sft2_5ce6c26e(result, expected, **options):
    """Check if image mode matches expected mode.

    Args:
        result: dict from getter with 'mode' key
        expected: rules dict, e.g. {"mode": "L"} for grayscale
    Returns:
        1.0 if mode matches, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger_qw35sft2_c98520.error('Result is not a dict: %s', result)
        return 0.0
    actual_mode = result.get('mode', '')
    if actual_mode == 'error':
        logger_qw35sft2_c98520.error('Getter returned error: %s', result.get('error'))
        return 0.0
    expected_mode = expected.get('mode', 'L')
    match = actual_mode == expected_mode
    logger_qw35sft2_c98520.info('Image mode check: actual=%s, expected=%s, match=%s', actual_mode, expected_mode, match)
    return 1.0 if match else 0.0

def check_gimp_layer_name__413739fbbe273da47f148920d82ae49b_qw35sft2_1ad7f5b9(result, expected, **options):
    """Check that a specific layer name exists in the XCF file."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    layer_names = result.get('layer_names', [])
    target = expected.get('layer_name', '')
    if not target:
        return 0.0
    return 1.0 if target in layer_names else 0.0

def check_gimp_mode_and_size__857686b0194013b75de5a057cd91e5ad_qw35sft2_7b9d3fdb(result, expected, **options):
    """Check image mode and dimensions with 0.5/0.5 partial credit.

    expected keys (passed as unwrapped rules dict):
      expected_mode  - e.g. 'P' for indexed, 'L' for grayscale, 'RGB'
      expected_width - integer pixel width
      expected_height - integer pixel height
    """
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_mode = expected.get('expected_mode', 'P')
    if result.get('mode') == expected_mode:
        score += 0.5
    ew = expected.get('expected_width')
    eh = expected.get('expected_height')
    if ew is not None and eh is not None:
        if result.get('width') == ew and result.get('height') == eh:
            score += 0.5
    return min(score, 1.0)

def check_image_list_created__1dbc4bff650a401c3959c27fb5e70015_qw35sft2_3181279e(result, expected, **options):
    """
    Returns 1.0 when ALL of:
      1. ~/Desktop/image_files.txt exists.
      2. Every non-empty line in the file ends with .png/.jpg/.jpeg (case-insensitive).
      3. The sorted file content exactly matches the sorted output of
         `find ~ -type f \\( -iname "*.png" -o -iname "*.jpg" -o -iname "*.jpeg" \\)`.
         This handles both the "no images" case (both empty → 1.0) and the
         normal case (file must list exactly the same paths as the reference find).
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('file_exists', False):
        return 0.0
    file_lines = result.get('file_lines', [])
    expected_lines = result.get('expected_lines', [])
    all_lines_valid = result.get('all_lines_valid', False)
    if file_lines and (not all_lines_valid):
        return 0.0
    if file_lines == expected_lines:
        return 1.0
    return 0.0

def check_transition_and_png__dd35705ee09fd488827a6afdd7cbab81_qw35sft2_86c528ba(result, expected, **options):
    """
    Partial-credit metric: 0.5 for PNG export, 0.5 for slide transition applied.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('png_exists', False):
        score += 0.5
    if result.get('has_transition', False):
        score += 0.5
    return score

def check_image_bottom_right__f123db274ed42f1e244c71177f056c09_qw35sft2_13a9b4f7(result, expected, **options):
    """Check that the image on Slide 2 is in the bottom-right quadrant."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    image_left = result.get('image_left')
    image_top = result.get('image_top')
    slide_width = result.get('slide_width')
    slide_height = result.get('slide_height')
    if image_left is None or image_top is None or slide_width is None or (slide_height is None):
        return 0.0
    if image_left >= slide_width / 2:
        score += 0.5
    if image_top >= slide_height / 2:
        score += 0.5
    return min(score, 1.0)

def check_image_on_right__2e5aa88aaf11fa4bda2a55acfcab9dc5_qw35sft2_52f1b487(result, expected, **options):
    """Check that the image on Slide 2 is on the right half of the slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    image_left = result.get('image_left')
    slide_width = result.get('slide_width')
    if image_left is None or slide_width is None:
        return 0.0
    threshold = expected.get('threshold', slide_width / 2)
    if image_left >= threshold:
        return 1.0
    return 0.0

def check_image_and_alignment__b349befcb19ab6c9d751d2833a5ceb8c_qw35sft2_e63d79de(result, expected, **options):
    """Partial credit: 0.5 for image inserted, 0.5 for image paragraph center-aligned."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    min_count = expected.get('min_image_count', 2)
    if result.get('image_count', 0) >= min_count:
        score += 0.5
    if result.get('has_centered_image') == expected.get('has_centered_image', True):
        score += 0.5
    return score

def check_image_and_pdf__cba3b2d7a3f4e1f51b9a84ac71e8e772_qw35sft2_7bb009ae(result, expected, **options):
    """Partial credit: 0.5 for image inserted in docx, 0.5 for PDF exported to desktop."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    min_count = expected.get('min_image_count', 2)
    if result.get('image_count', 0) >= min_count:
        score += 0.5
    if result.get('pdf_exists') == expected.get('pdf_exists', True):
        score += 0.5
    return score

def check_image_dimensions__b392913b04d99d3e7513b78f2dedf15d_qw35sft2_f9bcf566(result, expected, **options):
    """Verify exported image matches expected width and height."""
    if result.get('error') or not result.get('exists', False):
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    width_ok = result.get('width') == expected_width
    height_ok = result.get('height') == expected_height
    if width_ok and height_ok:
        return 1.0
    if width_ok or height_ok:
        return 0.5
    return 0.0

def check_image_compress_resize__8cc8ca1163a74ceeec7ab0fdb85713f3_qw35sft2_192625bb(result, expected, **options):
    """Partial credit: 0.5 for size under limit, 0.5 for both dimensions within max_dimension."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') and result.get('size_bytes', -1) < 0:
        return 0.0
    score = 0.0
    max_size = expected.get('max_size', 614400)
    max_dim = expected.get('max_dimension', 2000)
    size_bytes = result.get('size_bytes', -1)
    if 0 < size_bytes < max_size:
        score += 0.5
    width = result.get('width', -1)
    height = result.get('height', -1)
    if width > 0 and height > 0 and (width <= max_dim) and (height <= max_dim):
        score += 0.5
    return score

def check_image_resized_half__89f1952738bca9658039f047c420ba2d_qw35sft2_5c4c32f4(result, expected, **options):
    """Check if image dimensions are approximately half of original (160x255 from 320x510)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    width = result.get('width')
    height = result.get('height')
    if width is None or height is None:
        return 0.0
    expected_width = expected.get('expected_width', 160)
    expected_height = expected.get('expected_height', 255)
    tolerance = expected.get('tolerance', 10)
    score = 0.0
    if abs(width - expected_width) <= tolerance:
        score += 0.5
    if abs(height - expected_height) <= tolerance:
        score += 0.5
    return score

def check_jpeg_compressed__19b267768d7fcf0b434b3cca9c02b15d_qw35sft2_d7a03f53(result, expected, **options):
    """Check that the JPEG file exists and is under the max_size threshold."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or result.get('size_bytes', -1) < 0:
        return 0.0
    size_bytes = result.get('size_bytes', -1)
    max_size = expected.get('max_size', 409600)
    if size_bytes > 0 and size_bytes < max_size:
        return 1.0
    return 0.0

def check_jpg_move__7427978e92f6fc0e5652a4713261a5a8_qw35sft2_fe0d3627(result, expected, **options):
    """Check jpgs moved (present in cpjpg, absent from photos). Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_cpjpg = expected.get('expected_cpjpg_count', 4)
    expected_remaining = expected.get('expected_photos_remaining', 0)
    if result.get('cpjpg_count', 0) >= expected_cpjpg:
        score += 0.5
    if result.get('photos_jpg_remaining', -1) == expected_remaining:
        score += 0.5
    return score

def check_jpg_png_copy__ec3ddc36152b3d2fb4aee4f74f969e31_qw35sft2_14d06f7e(result, expected, **options):
    """Check that jpgs are in cpjpg and pngs are in cppng. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_jpg = expected.get('expected_jpg_count', 4)
    expected_png = expected.get('expected_png_count', 1)
    if result.get('jpg_count', 0) >= expected_jpg:
        score += 0.5
    if result.get('png_count', 0) >= expected_png:
        score += 0.5
    return score

def check_jpg_copy_with_count__58fb65f54a1152c1f3aecb552e15d932_qw35sft2_4feb954d(result, expected, **options):
    """Check jpgs copied and count file. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_jpg = expected.get('expected_jpg_count', 4)
    expected_count_str = expected.get('expected_count_str', '4')
    if result.get('jpg_count', 0) >= expected_jpg:
        score += 0.5
    count_val = result.get('count_file_value', '').strip()
    if count_val == expected_count_str:
        score += 0.5
    return score

def check_jpg_filelist__b885c2e3b122c9010091a38454018836_qw35sft2_eb1516da(result, expected, **options):
    """Check jpgs copied and filelist.txt has correct entries in alphabetical order. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_jpg = expected.get('expected_jpg_count', 4)
    required_entries = expected.get('required_entries', ['emnlp2023.jpg', 'hk_group_photo.jpg', 'hong-kong-china.jpg', 'monk252520thailand252520wat252520arun252520scaled25255B225255D.jpg'])
    if result.get('jpg_count', 0) >= expected_jpg:
        score += 0.5
    filelist = result.get('filelist_content', '')
    if filelist:
        lines = [l for l in filelist.splitlines() if l.strip()]
        all_present = all((entry in lines for entry in required_entries))
        is_sorted = lines == sorted(lines)
        if all_present and is_sorted:
            score += 0.5
    return score

def check_vacation_jpg_copy__91c6f86e45abe96db716a4e3d1072be2_qw35sft2_07c5ddf5(result, expected, **options):
    """Check vacation jpgs present and events jpg absent. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    listing = result.get('listing', '')
    jpg_count = result.get('jpg_count', 0)
    expected_count = expected.get('expected_jpg_count', 3)
    exclude_file = expected.get('exclude_file', 'emnlp2023.jpg')
    include_files = expected.get('include_files', ['hong-kong-china.jpg', 'hk_group_photo.jpg', 'monk252520thailand252520wat252520arun252520scaled25255B225255D.jpg'])
    present = all((f in listing for f in include_files))
    include_passed = jpg_count == expected_count and present
    if include_passed:
        score += 0.5
        if exclude_file not in listing:
            score += 0.5
    return score

def check_vlc_image_adjust__33b1c5d9e144110cf196db624efc6d81_qw35sft2_33445adf(result, expected, **options):
    """Return 1.0 iff the Image adjust filter is enabled in VLC."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('adjust_enabled') else 0.0
