"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_merged_csv__f43838d5904be041721c003e00fd36b5', 'check_csv_in_documents__a95ba08342a252657d47056ba988912e', 'check_csv_conversion__8b9d6ae5a27e52886aef51e39cc4df8a', 'check_csv_export__02188a8f533477bfdcad341f72281666', 'check_csv_filtered_a__2fadb675c56bf8de284fee43324487b2_qw35sft2_2bdb074f', 'check_csv_space_merged__00d80996b3209df110a0b1e69fc0ab31_qw35sft2_fd4c6711']

def check_merged_csv__f43838d5904be041721c003e00fd36b5(result, expected, **options):
    """Check if merged CSV has correct header and first few full name rows."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    if not lines:
        return 0.0
    score = 0.0
    expected_header = expected.get('header', '')
    expected_rows = expected.get('expected_rows', [])
    if lines and expected_header.lower() in lines[0].lower():
        score += 0.2
    if len(expected_rows) > 0:
        per_row = 0.8 / len(expected_rows)
        for (i, er) in enumerate(expected_rows):
            if i + 1 < len(lines):
                actual = lines[i + 1].strip().strip('"').strip("'")
                expected_val = er.strip()
                if expected_val.lower() == actual.lower():
                    score += per_row
                elif expected_val.lower() in actual.lower():
                    score += per_row * 0.5
    return min(score, 1.0)

def check_csv_in_documents__a95ba08342a252657d47056ba988912e(result, expected, **options):
    """Check ODS to CSV conversion saved to Documents directory."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('used_terminal') == 'use terminal':
        score += 0.5
    expected_min_lines = expected.get('expected_min_lines', 5000)
    if result.get('file_exists') and result.get('has_valid_header'):
        line_count = result.get('line_count', 0)
        if line_count >= expected_min_lines:
            score += 0.5
        elif line_count > 1:
            score += 0.25
    return min(score, 1.0)

def check_csv_conversion__8b9d6ae5a27e52886aef51e39cc4df8a(result, expected, **options):
    """Check if CSV file was created with correct header and first few values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    if not lines:
        return 0.0
    score = 0.0
    expected_header = expected.get('header', '')
    expected_values = expected.get('first_values', [])
    if lines and expected_header.lower() in lines[0].lower():
        score += 0.3
    if len(expected_values) > 0:
        per_value = 0.7 / len(expected_values)
        for (i, ev) in enumerate(expected_values):
            if i + 1 < len(lines) and ev.lower() in lines[i + 1].lower():
                score += per_value
    return min(score, 1.0)

def check_csv_export__02188a8f533477bfdcad341f72281666(result, expected, **options):
    """Check CSV export has correct structure. Partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    if result.get('exists'):
        score += 0.4
    expected_headers = rules.get('expected_headers', [])
    actual_header = result.get('header', [])
    if expected_headers:
        matched = 0
        for eh in expected_headers:
            for ah in actual_header:
                if eh.lower() in ah.lower():
                    matched += 1
                    break
        if len(expected_headers) > 0:
            header_score = matched / len(expected_headers)
            score += 0.3 * header_score
    min_rows = rules.get('min_data_rows', 40)
    actual_rows = result.get('data_rows', 0)
    if actual_rows >= min_rows:
        score += 0.3
    return min(score, 1.0)

def check_csv_filtered_a__2fadb675c56bf8de284fee43324487b2_qw35sft2_2bdb074f(result, expected, **options):
    """Check that filtered.csv contains only rows where last name starts with 'A'."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    data_count = result.get('data_row_count', 0)
    expected_count = rules.get('expected_data_row_count', 400)
    if data_count > 0:
        score += 0.33
    if data_count >= expected_count * 0.95 and data_count <= expected_count * 1.05:
        score += 0.34
    if result.get('all_last_names_start_with_a'):
        score += 0.33
    return min(score, 1.0)

def check_csv_space_merged__00d80996b3209df110a0b1e69fc0ab31_qw35sft2_fd4c6711(result, expected, **options):
    """Check that the output CSV uses space-joined full names (single column format)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    data_count = result.get('data_row_count', 0)
    expected_count = rules.get('expected_data_row_count', 5000)
    if data_count >= expected_count * 0.95:
        score += 0.33
    if result.get('has_space_format') and (not result.get('has_tab')):
        score += 0.34
    first_data = result.get('first_data_line', '')
    expected_first = rules.get('expected_first_data_line', 'Dulce Abril')
    if first_data.strip() == expected_first.strip():
        score += 0.33
    return min(score, 1.0)
