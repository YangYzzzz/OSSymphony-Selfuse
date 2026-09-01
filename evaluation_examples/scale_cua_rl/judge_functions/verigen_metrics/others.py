"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_row_scores__dfb627af00a19217031ee997d3cc240d', 'check_sorted_columns__e0472552b507a6465c494b007b9b0305', 'check_earliest_entry__53a0fb3c44861b74092bea203bb878cc', 'check_descending_sort__f8df3778414dcdc67e5f887268cde89d', 'check_sar_memory_report__a8cfa9828c8ad4ee330aec94e4adaf4c', 'check_total_row__bf5b353ae5739d8c5c981c28dcf67581', 'check_year_values__5d33fe5f963081b2f481ac5f73c3ccbc', 'check_dir_structure__f789180956e4b0547b5929aa18a0352e', 'check_budget_reservation_info__a4a15f9cf85bf343423f5b15bc460e50', 'check_heading2_alignment__131e78d644cdd4aebe878331238d0634', 'check_main_caption__909b46690f7fc746deb97149427ce266', 'check_name_and_count__009b0a2429b483dc78b74a0168744c1a', 'check_pass_fail_held__d026dfbc396400238630feee90a0bfc7', 'check_frame_extract__5e444f1b525a6375eead73fcf3fce163', 'check_user_homedir__df678257c48d0ebe8f2e03968592b783', 'check_compose_cc__457e68f7ef86811b4470193bf0ca1e7e', 'check_row_hidden__18d397a77a413967bac86e92a1394add', 'check_navigator_visible__a0fa1d9cceca2f2896bb79deb4b8d794', 'check_user_password_home__6fbd0fb5ef1077b146e8ceece522a6b4', 'check_picture_bottom_left__77f9e57f0af70c9666efd508929f72c3', 'check_row_scores__ba349a57c83f48325e3b8e8c17891359', 'check_background_color__0fc932ce1ad9456b6d290262a9551736', 'check_dir_move__a7ceb843bee94531ff589f6c5602590a', 'check_column_deleted__51f83653b8f37d64288aaccfa8443c8f', 'check_cs_data__6d4ad198dfe8ff0251a084115fc1444c', 'check_chapter_list__a4d1b434aaedff8ce21d5e5d82eac67e', 'check_manifest_structure__783369d0e7e20c29051ab5abdc75ebb3', 'check_orientation__060b6ae42d4d041f46057b0dff06cb53', 'check_station_removed__dc2cb3918367317ad49987158bb0de7d', 'check_venv__3936cc1f2257a01f73d410787af3a332', 'check_multi_replacement__43d92268ad8323e8d3f0af9cc0d7b2ce', 'check_optimized_bubblesort__8aa2e8e2be1ca5d610e977b9c78060b2', 'check_do_not_track__956b75b2b213e77466e1066455a2b384', 'check_verified_invoices__445a1e82d5445803b67308e02fdc3544', 'check_enabled_experiments_contains__1656c6075ec2ba63a026b59d9c69a38f', 'check_notification_dual_settings__0a0313fec36cad178ad94d7913e78201', 'check_summary_rep_totals__8a3eb44c98a42fd6dcd94122bb8b7a9a', 'check_period_rate_col__02a40d6cc593a8648c91f6def88b5342', 'check_blue_color__be59b6debe1fbee8fc53696bdb2273b8', 'check_enabled_experiments_subset__197c5f21c248c0028a67e57d9193addd', 'check_cumulative_col__6bc8ae788974152807b9723497a26da7', 'check_contrast_increase__a61095a992c03ac5e17b282d64e6f2d0', 'check_top_performers__85190e1339a2cb1782e74f01270bc172', 'check_invoice_total__aa615bdde918e746d9793d851a2bb4f8', 'check_do_not_track__8d6140655fd3ffdbb0f2973a3475732a', 'check_do_not_track__baa6ff2cd46140d8d50ba44da41af676', 'check_sum_row__99d61f333e953d709c5fe5f221bfb63e', 'check_compose_recipients__9a5e3321e7daf20a2f35cb35407bc19e', 'check_sales_by_rep__5c7590b2fce2859cb43c07916cdde25c', 'check_profit_column__11a8fd5475e2aba6ae95fda21befab56', 'check_rename_copy__25257c2d579c2a114cd376d84d920386', 'check_do_not_track__7317fcab347872285f8a0f137be4010c', 'check_novel_info__6a9459e3699067e60bf00ce031fcbd18', 'check_net_income_column__bef74e7adbdaaf2c559cbe5b3bf80940', 'check_contact_export__091acaa52f197bd2f2171a56782c0841', 'check_row_values__4551fb3f098ec9fe882ac83568cd431f', 'check_columns_hidden__a4f4a609a17cec62166f8b2d532de58f', 'check_row_values__0cd3deb0e30597bee67bd1ee61842fcf', 'check_column_sorted_ascending__0f45568bd733d9527d872ad7c053c2c8', 'check_import_exists__d8b2f0ccd335f9a3a3203363805b5ac0', 'check_badge_deleted__6b4bf280e2d30f49e851d7a1a643ef02', 'check_merged_novel__2abe2760620ddfe22faad3f384e1cae4', 'check_dir_structure__95c704f58705a26de50af292b2eb5059', 'check_model_names__55e51bced75ba3e5e08e1be190df564e', 'check_window_size__e47348d3e93e8a085dab2f1753ad88d5', 'check_song_list_content__09209b5893b5a25bd669841ecee9c784', 'check_zip_contents__b6394280513f3b0b8d3eaa6ea7ea0992', 'check_compose_subject__df6702fb4aa059fb4a86e567f6f9f62d', 'check_grayscale__71dd5cb242259598df523724839bf008', 'check_dir_ownership__425163a454302c61dc411d05ac9420b7', 'check_do_not_track__fef997178e549d845fc31b37602ff4d0', 'check_descending_order__40424b53798c2f42a45c9b5d1f78c9a0', 'check_delimiter_replaced__0a313c15fc4fb10c48c753cea47e4527', 'check_sorted_names__56c311ddff224f797b6614d5151a3c35', 'check_mpstat_report__28ae802c5923a6e7ecfebd0e5cc59702', 'check_sorted_values__0c918bbe7b80a9e03be9fde4824c7824', 'check_title_style_bg__d816da51e2714a5b86252b84b6e77125', 'check_title_added__2130fa798c3c78ff755136f88d8b3efd', 'check_fps_value__6d1f5614df32aa00fba65a5a2ca43cb2', 'check_title_centered__f845858287fa4f02431e180c3025ad53', 'check_strikethrough__a345b8271c8198e8c95ec913d78c781f', 'check_do_not_track_enabled__2a127fa9e4cd79198d970c24eb95cab0', 'check_word_counts__4ca76e68b8033d8d7fa4f889ee4aceef', 'check_booking_form_state__cf1c94ef66b2087660e9bbd9fc7c9d7e', 'check_screen_lock_settings__39d04f914d2dbd8627ddcd7f791e785b', 'check_zone3_row_totals__c68f755a470f57d0018a142ca4ed4f38', 'check_chart_exists__9affe8a59575f48a48ea3570040f007c', 'check_column_sorted_descending__3201def6075c63c69d0100c10729d61d', 'check_range_values__160ec07bd1e641be9b31cf92f3848b19', 'check_grocery_entry__00c370b29bbdeefb624743c61a277c9e', 'check_title_bg__c2ac2f94e6c641c5a5c2e733ff74ade1', 'check_model_count__204b2db58c306983ce47277dbc39788e', 'check_empty_count__4128b5119649f756a2c7d9fc4321d25f', 'check_total_row__6a822179da68f351e1e0a8c5b6636775', 'check_conference_continents__a9d0754cd398d25ba1c2a300dbecfcf0', 'check_manifest_structure__ab85f9bbe05040b98eb1a964cbee7866', 'check_title_format__110ae328b957703cc4b8db2a2b04df5c', 'check_manifest_structure__61a5e99298254bc3ea7017f23352ee68', 'check_row_values__8c5594571c596e8d15ba168f23066d41', 'check_repair_entry__c6556bc153411ec2613fae626891d4a6', 'check_column_values__6d4300c0b784865a64e946f955774003', 'check_heading_alignment__e69f89d4207a469ba2d4f747df580e7b', 'check_units_by_product__8ec7487c1752ceb36e17eb91a13203df', 'check_uncomment__cd3236440d08e0a8882e8811427d01d3', 'check_vlookup_prices__40371a0bf41fabc605f9425ce4121622', 'check_upper_titles__6d702aff56732f3b1cb6a80e7f740e70', 'check_scaling_factor__d5d03d50dfc17bfe4f646f3d3c9160fb', 'check_tally_entry__89c74c0ac9693815f604c967708f8f82', 'check_power_settings_combined__e637b95c625a88b3bb2cee1c77022b86', 'check_title_alignment__4f408a61496e9bff25fa5074f0f2e380', 'check_import_extraction__c191ddd519b4043306a4effe99ae0aac', 'check_column_values__4065e6a4b9dbccc0390b8c5a38e2d1b7', 'check_settings_speed_blocksize__4b546ac4b7d86d470dd227a5c08930e6', 'check_sorted_sales__38a9aeefb36e58a6a0dcece70b9af3ed', 'check_model_list_unordered__68e1e5b5d5b054d1e9ca98b707844de8', 'check_half_rate_col__fd58351f5fafe47d305d3595c5202dcd', 'check_rows_hidden__cb276526eec2ec94231e7dedd20b5069', 'check_fv_column__5c1e2dfe8bffbdb20d88aa7499677c46', 'check_utils_function__597b85b971b86604637b8eea77967812', 'check_column_sums__427c1b1ebceb5c2ead27edb218ee7bf9', 'check_year_values__dcb80f39266c26be10a70749c6e4e4f7', 'check_do_not_track_enabled__66408ac42f837225d56e2756a129da0a', 'check_net_income_column__bc1bb048c0eeb30a1fa5f9604d4f0a11', 'check_code_stats__7636602600807a890f06c684e5901e16', 'check_names_and_city__95be0722e22941aff21eff3d680c0c01', 'check_move_symlinks__2a402a2f74167209d0ebc1db72b76fd8', 'check_lower_titles__56ea06753821989747521678ff3594b3', 'check_row_scores__1700c01d5ad41ab8401de0bafb801990', 'check_column_sorted__34bc68338cfc2efa5ed7f80b3fe4afc3', 'check_title_color__6b28d17bfe07bb684917211f85bf497c', 'check_inbox_state__8e4a0a40e4d045486197a43a70e9b321', 'check_zone2_col_totals__dfa00147f689f96c837b3a22754c016b', 'check_column_values__2b83d943849f7c327e47bb940fe75967', 'check_ascending_order__5f28938c5d9137ac59ddda20f23ca5f2', 'check_bibtex_entry__dd385e64b1e5f06bf4b4c074020cb6a0', 'check_conference_countries__70ea2799be65bb426ef2e5f3f76ece43', 'check_column_values__312ae9a7a4fdef961589d00f7ec76217', 'check_row_values__b46e00d9bcf11154ef195dbbf421e737', 'check_function_rename__bef206759cc52e3fdba71a4f0e77e8b3', 'check_horizontal_flip__30cadb0ec7eca28df693cf924ab88301', 'check_row_values__767f7417cb0fc080051f9d7103e674cb', 'check_no_pictures__48004e0f145da635f1d1ef72ec0c75ce', 'check_notes_bg__7bae80483f423bfa3b3bbb7f9e9eafdb', 'check_color_and_size__043ea3777ab786b45abdbd440171ffc3', 'check_dir_exists__93ffa0d39e3e07dab53be514a31da362', 'check_gemini_response__ae954eaf7f0ec05c9c2918c4eec73bcc', 'check_number_format__092734e9cc221c6f71e099e209498fb4', 'check_currency_format__6f34d4ab7744f84c76fe24f84ac55b09', 'check_padded_ids__598fab012ada3ac788280b75e5c01391', 'check_sorted_ascending__85a76f9ca092376f5f612fc019e6af39', 'check_highest_price__4c5e4ba8ba69ff54d1e2443439646d44', 'check_columns_filled_and_sum__f0c5ce9b9ab74cdf70d3ea75ec500ce3', 'check_numbered_notes__5e14eceba731d7ebe21b244525a2edc1', 'check_zone1_row_totals__b15b7a7665c71e3a9519658d5a13d57f', 'check_backup_move__1a7983e090de51b7f1a1b54e751935d9', 'check_sidebar_hidden__c002e9fc79af5f293a68e76c7f778a2e', 'check_vtt_conversion__9ac6e58c407c5a10e99c8a824933ab2f', 'check_aws_total__d9fdc5409ef8a872907b1ddc3db9cd54', 'check_row_count__79ff096ee36ebe0503f424e88a332ca2', 'check_row_values__19698d046d0da22c68c3bfcbfbd57b88', 'check_sharper__dd81583a369d4d4452910cdac2a642dc', 'check_def_extraction__9f211d2fef2982f0959041bc806f53c4', 'check_restaurant_names__a68724700bb60550a33fa328934ea4ee', 'check_column_sorted_descending__542fc695553e3e27095b27327032df35', 'check_iostat_report__d175f88bb1991b0561f10b03c2a108ce', 'check_ordered_list__0762adbae4ad235c185f47096cf64c91', 'check_sorted_order__4841f866c3a0d68d763e22faff3652cb', 'check_vcard_count__36ae571726da697a61321c69b6ca11f3', 'check_eml_backup__b0f7174902a45f1ac3121d8b78a62ff1', 'check_profit_margin_col__60b13be5a7ccd193b76ecf097a198265', 'check_multi_column_values__9c6205dd9e0fb3145a3021bfaf014d23', 'check_sorted_invoices__5d7cb95eaab32c834b0336c7360277e1', 'check_paper_titles__db9975908e0c2c40a4467f2add6ed446', 'check_content_align_title_color__63833ad040c5e27e84bd66675fcf45fe', 'check_settings_screen_dims__29161befcd93aa27778463661523fc23', 'check_safe_browsing_enabled__949ee9014650e2c21e76ce453947e00b_qw35sft2_7e58fdfd', 'check_do_not_track__59ed9e9855d0703d98353105912319d0_qw35sft2_eee1b9e4', 'check_delta_miles_checkbox__f9b3359b12066a252ed8849a9876d835_qw35sft2_e03606bf', 'check_delta_trip_type__45facd29054b2deec3804036a06322c9_qw35sft2_35e971f8', 'check_sb_and_dnt__497c76e8a817211059a0ca1d7eafcbe5_qw35sft2_83f45ad6', 'check_startup_opens_google__db8335a9df035ba863791bb757e40ea6_qw35sft2_08f31e16', 'check_appearance_and_dnt__197c5f21c248c0028a67e57d9193addd_qw35sft2_2650bcc9', 'check_lang_and_dnt__45fdf66a8d514412f8287b18b125d4f9_qw35sft2_4a99185c', 'check_do_not_track_enabled__ecfeb6992b96a90fa421f7b2969c282f_qw35sft2_e39772e4', 'check_recreation_title__6368f69cb89cbc42bff053c412aec363_qw35sft2_860c41ff', 'check_delta_origin__607589a24fcd7d8000928f466f0dd577_qw35sft2_c52f3ac4', 'check_delta_destination_field__696c0a5477f795f89ba4614ef9b530f9_qw35sft2_7ce9bd2d', 'check_history_and_dnt__43a7a4d8fb3eedcedaf50f5a89a2ea93_qw35sft2_2f3f6db4', 'check_enhanced_protection__ba99e87e44fe809095fa3ce0e0316f1d_qw35sft2_938306e7', 'check_contrast_and_saturation__b674a0dc75a50659a17535e66cb88a78_qw35sft2_546b48af', 'check_frame_at_3s__81411ca26f8f6e6edae115d08d9ab086_qw35sft2_01cab48d', 'check_brightness_increased__1dd8caa614c12c2ce7c1b58be9a02c11_qw35sft2_4166e7d3', 'check_grayscale_and_contrast__a80023f11b098859420cd69b0067e271_qw35sft2_8b42bb15', 'check_brightness_increased__53b38a939b32dab621a5954331033a3c_qw35sft2_ab8758a4', 'check_scale_and_contrast__4ce9d4657814d8ed93df9fa8b61db4cf_qw35sft2_13cdc41c', 'check_three_brightness_increased__b3992f4af3f3949be116d5f6a9289f79_qw35sft2_4c2ee848', 'check_freeze_panes__10368a4827622b87edd2c56b2f249e0f_qw35sft2_6266b504', 'check_salesrep_jan_total__5238964d8e657c42c8a059231010b871_qw35sft2_f69af80e', 'check_level_secondary__f558bfbb444b8f5736dbe2588385d5fe_qw35sft2_dd3e1b01', 'check_income_gross_with_total__9bcff5517fb7493e61233c0a423569f9_qw35sft2_ea0fbb59', 'check_employee_ages_avg__151640eb3d43a7eb22c551550e103953_qw35sft2_fc1f68c6', 'check_vlookup_f2_single__dba068db6ac4d383c892aae7e4d7fc62_qw35sft2_c217ff4b', 'check_sort_total_row__2cdc28af4ef71e22984f659197334a50_qw35sft2_aa12df02', 'check_total_label_and_jan__4e23e64ce56ffa608b1e8fd2866f5ae0_qw35sft2_77da4f1d', 'check_first_col_not_empty__b3e437bf5bf6ec56d9f1c4d46601f3ff_qw35sft2_8a9fd9e9', 'check_monthly_totals__aff1e90d6f0bac6579c8d1279e6c6c3e_qw35sft2_a24f5cbb', 'check_weekly_sales_profit_total_row__50f7148385ea6d504119cd2e40f7d1be_qw35sft2_2b548c12', 'check_employee_split_sorted__5d1c2e4438b4aac3fb5db710d29300f5_qw35sft2_ffeb6422', 'check_c1_format_d1_fix__3be2ab82396c7145b9caded82bed3999_qw35sft2_d313d34b', 'check_period_rate_max_in_d1__14298c6976685b9e8d10b60e8490521a_qw35sft2_e583f219', 'check_sort_and_sum__8b987ab7e49940ce5b75cd3329aaf643_qw35sft2_0cb12aa2', 'check_maturity_sorted__6a0b6ecdb454d8389c50fe9e251b1580_qw35sft2_806b62f8', 'check_passfail_and_validation__1231c66fffc58f9761299238f8444bbb_qw35sft2_b3126219', 'check_salesrep_label__d6d9e5e1d2ca3d967d1969fd1a1e0c2c_qw35sft2_7b6ca218', 'check_seqno_and_total_row__f9ba38ec0e00ec28c9c0ef5057c79916_qw35sft2_0d24a241', 'check_employee_ages__add589b113a1be1eec737c1ab22661fd_qw35sft2_2484c826', 'check_vlookup_f2_f4_rows__22ff441442432fb67bd17ef3eb97a292_qw35sft2_eb420e13', 'check_income_net_sales_gross__c6e3aa9fd1cbfd34b4473b9e8f31e349_qw35sft2_80a50e1c', 'check_level_primary_range__74bebda311b0ec1306e448288c71ed13_qw35sft2_10d2d467', 'check_total_row_state__5ac69261c8fb294b16f18049ece06ce9_qw35sft2_8f27b526', 'check_single_row_hidden__93bd58adf830a704a8880938e35d28c4_qw35sft2_d3296ba3', 'check_weekly_sales_sorted_by_profit__99a00a677b41fb8db78ec79d5381f0b2_qw35sft2_fc61000e', 'check_feb_max__663a803da15ba1a41afdeb0114613f5e_qw35sft2_31440d0b', 'check_d1_usd_fix__62d484d2d5eca5e0b713e4dbdcc095e8_qw35sft2_dcdf4e3f', 'check_ramp_accel_diff__e64652ba3bf9a74b525e231bf3f391d2_qw35sft2_dbfad9cd', 'check_period_rate_sum_c26__fd409a83b18f09cf45613c5173377d4d_qw35sft2_cb0acd0d', 'check_sort_and_sum__7b14f7b3fd242abdf7c663f7ccb97562_qw35sft2_d1f6f816', 'check_sort_and_sumif__bda50c4eee48b4e3a15799921e84da62_qw35sft2_5d609264', 'check_sales_total__c61f5322a84515967a3c4cc16d8ae654_qw35sft2_4f7f37c0', 'check_maturity_total__461252d0c8044977d30ecc240f2dc9cd_qw35sft2_677d3c70', 'check_passfail_and_count__9a44e82214177e638fd33343975b4139_qw35sft2_22da825c', 'check_row_freeze__ae87df329e6b7ba3f255d9a2acdfeda6_qw35sft2_0ffacc85', 'check_seqno_and_sales_sum__00713006369c01f94f764c7eb6008543_qw35sft2_1ee1692c', 'check_salesrep_last3_totals__ce77d41e3a7b3efa48e4655c22c8aa27_qw35sft2_5bfb12c9', 'check_student_c3_c4__87a420fb18a4347285a9615e2d7a9d87_qw35sft2_1102acc6', 'check_vlookup_f2_f12_all__202132c1158925d79d1ba222174d8f66_qw35sft2_934b09f4', 'check_income_net_sales_and_cost__43eee44b990a73c4b5a6e2bc138833a7_qw35sft2_b852ddbd', 'check_save_and_transition__5e93f9ddd16cd129b875ad7386e8a5e5_qw35sft2_17efed89', 'check_presenter_console_disabled__c6ff7e494d134a42d6e5c420e00bf4fd_qw35sft2_502d7a7b', 'check_stretch_and_no_picture3__bd30ea01dba8f440569b065c94e4f32f_qw35sft2_1c1836e7', 'check_strikethrough_and_transition__b374fc122f888120a74e40bee1da5133_qw35sft2_09c1ad4a', 'check_stretch_and_fade_transition__676dd2418e7e3a84c9083aa5ad9f76e2_qw35sft2_15654b86', 'check_picture_heights__eff65ebb8ed6d102db81f31da1a823ed_qw35sft2_de0e2541', 'check_strikethrough_multi__7f3881fcd58bcba327d08f20ed190dcc_qw35sft2_e7c6ddca', 'check_titlecase_and_center__c181651f58b43b69af608583a0523891_qw35sft2_81f55bca', 'check_word_colors__336f1c4245e680e1aace133fb243059a_qw35sft2_e8073d66', 'check_first_two_double__c985596bc3953d61c7d274fc30ba4960_qw35sft2_b823e19e', 'check_word_colors__d89e455c3ca4dbc17e773411cf8f66df_qw35sft2_d3435f5f', 'check_title_alignment__6435e9b69d0bbcb863c89964e3084688_qw35sft2_ed3834f4', 'check_word_colors__e97799f0add4268fcc31cdd1e95ac277_qw35sft2_e161288d', 'check_mixed_spacing__ffe17a8bd50575f09eec01a68fed60c6_qw35sft2_e0536b7e', 'check_titlecase_and_doublespace__6e2b5ade87cb088089f67dc5070eb77b_qw35sft2_867c1751', 'check_word_colors__3d9f327a66b9025c03e34cafe8c88fe3_qw35sft2_fcfc98f6', 'check_settings_window_size__a8bc55c18275bf61543bfdfa6bb3d447_qw35sft2_068a7c7a', 'check_dir_exists__e66cc6b0ec6479dc303b558a9cc72642_qw35sft2_43067696', 'check_res_txt_descending__da27b9fa71ba67953d0b9f649a3197fa_qw35sft2_7fd0b028', 'check_invoiceGES_in_problematic__cea86e2d544f9d987b2adcf034f36c4a_qw35sft2_6981b395', 'check_2017_2018_cities__1272197382489180193ce8b35860bb11_qw35sft2_f2b45d1d', 'check_small_compression__98edc79a225cc1dd6b5d8a23724241ab_qw35sft2_cc37c421', 'check_answers_tests23__ed93db9128130b7a782216ac4507a0ca_qw35sft2_80825339', 'check_sar_disk_report__fc8ac0599dbb02ca56427b3ae265c796_qw35sft2_a50d1235', 'check_txt_multihop_gemini__61501dfdaa17fec1d3c116cb7f744d5c_qw35sft2_3940d7cb', 'check_imdb_top250_nav__7c977607bff83aa73fb51b31d1f513df_qw35sft2_0c77bee4', 'check_dir_exists__e103d8977cd2fefeec7972ff8169e36d_qw35sft2_b561c678', 'check_settings_snake_size__85e39531da2c848a1afaf10e25ba3dad_qw35sft2_a4edccc9', 'check_invoiceTII_in_problematic__7ecb63a4609903641df39f4754b7248f_qw35sft2_b90e9695', 'check_book_copied_to_documents__55d3e622f74c4c6d9051e6e358a64ecd_qw35sft2_b3db2c69', 'check_bubblesort_impl_and_output__8c65eee451a1bb106a548123af80ea08_qw35sft2_3002e435', 'check_ext_and_resize__3001e2e091ccfaef8fccec23c5f15501_qw35sft2_ad209965', 'check_sampled_conf_cities__fb22e54de099b00f4d6a127b12703f11_qw35sft2_b8b70a1e', 'check_tally_book_amount_sum__289ae9cd25fb72b631d3d5c97c4b9f82_qw35sft2_fdfcb7c8', 'check_paper03_and_year__d508d903f596b0ecb03360bdda6892d4_qw35sft2_f08f2902', 'check_goodbye_world__f65942d684aa73fa8a727494a6335877_qw35sft2_c83e2311', 'check_sar_cpu_report__5dfc721833abb9e182bb18218c19d630_qw35sft2_1e2c5a0d', 'check_answers_and_count__25c4c78ad235f6f7dec2a3ce1f7f2c6f_qw35sft2_30c1ab5f', 'check_settings_fps__751d0a9ef75d8ec2c4c27e57248ad37c_qw35sft2_daa09262', 'check_invoice243729_in_problematic__da5edbd81a2d596fb7eb6ef7b52c2151_qw35sft2_2c597180', 'check_clipboard_path__241466d2f7566981ed09035f56716bb4_qw35sft2_6fab5826', 'check_conda_install__d5abdc9092a51546a82222c972d5edf6_qw35sft2_d9f2a9fa', 'check_icml_cities__e73f742a9bd3e72f68e6861886fdd0c7_qw35sft2_25ed0820', 'check_paper01_and_count__101069bcaae5b37861b743495ae4859d_qw35sft2_03bcee0a', 'check_res_txt_sorted__68c8947721f1221211aba42cd6d1735d_qw35sft2_c967d941', 'check_tally_book_last_row__97f6f4530ffb08062fee8d0568f59984_qw35sft2_d1e6ad53', 'check_copy_move_state__ec373221258728f3163857f2bdfd353a_qw35sft2_a2d7b604', 'check_clock_dual_settings__443802668a0bfacbf0757782ecb2ad43_qw35sft2_e19209e8', 'check_screen_blank_timeout__40046a4f8d5906bbddcb5abf1fee6b07_qw35sft2_1cac0dec', 'check_screen_lock_enabled__74e051f75a4f9998d114f6407d55ddd5_qw35sft2_c77e69c9', 'check_volume_level__d6b3f95410b1cde5a46558c01aaafc46_qw35sft2_4d0db765', 'check_utc0_24h_clock__c55229b53381e6587f722d00658cbd02_qw35sft2_68eb61e3', 'check_copy_4dirs__4f3015d5aa890d8fea505e363d3f7aff_qw35sft2_13a61b9c', 'check_rename_and_sibling__9cc29f7146497c596e333ae0ebaf5848_qw35sft2_964736b0', 'check_notif_clock__accf3abf7a1d21f7f68eefd98414a494_qw35sft2_cd16c09b', 'check_power_dim_and_blank__4a865a2b73390cd1be0bc5f929d03ff5_qw35sft2_fc5c977b', 'check_screen_lock_with_delay__8b9cedf68783a02fb2242a963784d1c2_qw35sft2_5da34c49', 'check_php_stats__a63cdece1fd90a60a3e737af233f1764_qw35sft2_aba72e41', 'check_move_failed_notebooks__96695aa3604bded58ed2ec7f3c72c8de_qw35sft2_199fbbbb', 'check_output_and_backup__9fff5b0d3da3f15287b691b4b1944f6b_qw35sft2_939b4b84', 'check_rename_and_move__eb200f79b699eb5d233b5379c7e1ea0e_qw35sft2_8153b443', 'check_volume_level__2ac92884cbfeefbd5a2d6929c20839d7_qw35sft2_5e6ad00c', 'check_screen_blank_never__3696e962dd7f0719e2ddde3289beabec_qw35sft2_89043873', 'check_restore_and_copy__f767d224a0cdea04b11256d9cffabd41_qw35sft2_49933203', 'check_copy_and_count_fails__8f62d3b2706e8f98a78e22aea189146a_qw35sft2_24b70f63', 'check_notif_dnd_lockscreen__7a1296a554f735dccd7fa86e30d68dbd_qw35sft2_9d57f044', 'check_utc0_ntp_disabled__397a38a7528d157681d1e57316ff4073_qw35sft2_008be2a1', 'check_screen_lock_with_delay__b5802fb719515ed2e2b308492efceb55_qw35sft2_3bd8d49f', 'check_subdir_split_permissions__31b8299b7936d108e594f17a62518506_qw35sft2_16fc09a7', 'check_volume_settings__1201dde5af3c06b515ed16eae4fad685_qw35sft2_279ea4a6', 'check_restore_and_rename__ee992b5c4f2de4a7b0ef0f22ec1a67b1_qw35sft2_a576718e', 'check_notif_screenlock__0ee0cafb23f5579534516156a3b19db3_qw35sft2_0abda510', 'check_power_triple_state__570bc607f27286d53f7b1a8eaf2af624_qw35sft2_ef529bf8', 'check_screen_lock_with_delay__5e9c7d265f840de83da743ff562a1457_qw35sft2_7bb5e860', 'check_output_and_count__bb3776631f95d3c4e2bcd96582b451ff_qw35sft2_a7da54b2', 'check_restore_and_wallpaper__354747c46b7fabc52ac1f7b0a9e657e8_qw35sft2_0291f617', 'check_notif_sounds__139fa3103c161a5284bcac4b5dc4f632_qw35sft2_2a8c221a', 'check_power_dim_and_suspend__6cd1c7feeea3fc52512c1552e8d5b3b2_qw35sft2_c03caca2', 'check_copy_and_rename__44eff79b79182e6eeb0c7debb1bc9f8c_qw35sft2_0bd73cb5', 'check_archive_content__7b7da38700f679a895053728c5a7d35b_qw35sft2_c00fd3bb', 'check_two_dir_notebook_split__e71563318b46b13cadfcb76e6b9ce246_qw35sft2_e2a4d5be', 'check_accessibility_two_toggles__8ee8727a5fa2032ce96fa89a79ac9fab_qw35sft2_da39f4c7', 'check_eml_backup__d81f9d153146f886b82c4131d09d1ccb_qw35sft2_7f674302', 'check_bills_starred_and_important__ed9097e7c75920408b4536024b8da2a4_qw35sft2_53fd849d', 'check_bills_triple_combo__d21b0cc63afcca51c02fdd16b5d862e0_qw35sft2_a12fd2c7', 'check_eml_count__2731b9abd5cfbad9ed4df8aae737addc_qw35sft2_70f9d6b7', 'check_bills_starred_unread__8e8eeb1588f1109e98ccb02a0faa787c_qw35sft2_e7b5f644', 'check_eml_subject__dfeb48225188ee18fb4de9d6f0048829_qw35sft2_631026b2', 'check_bills_starred_and_receipts__0b7b1b3c91e02a8222608f02952fc214_qw35sft2_aa546db1', 'check_bills_starred_and_filter__44c5673d772aeaf1fd4e6dcc8a32111d_qw35sft2_624b5785', 'check_play_and_exit__ffd6f40d9680fbe09c9ab1f8d445d696_qw35sft2_3a30b889', 'check_global_keys_play_stop__0ab6bfa84809c927c91a6a1387eadc16_qw35sft2_1540c5d9', 'check_global_key_next__d8953e567823b8ea130733a22ba09b75_qw35sft2_1b503325', 'check_global_key_vol_up__d76c6ff727ecd69650cf5e0e8eb5e19a_qw35sft2_b23ce3cf', 'check_play_pause_and_recording__0ed5ebba8c137cdb9537f1717121259d_qw35sft2_4ed02195', 'check_global_keys_play_next__2b7ee2d73c0169078c29f1690d78cdca_qw35sft2_1bb2ac84', 'check_ext_and_multi_settings__afc2fb3b68b53df8d476e05f07933aff_qw35sft2_4a8a2864', 'check_dual_replace__ae0b3ab6740bff99e232fccee3b36aad_qw35sft2_358b4d9f', 'check_dual_replace__004ad4b9f1b7c0d57d366b67195586e3_qw35sft2_0f6dfe81', 'check_ext_and_wordwrap__324a83eafb9ff6ed07fafe7e199af4d5_qw35sft2_cd76830c', 'check_randint_syntax__2c85ba49c59e818794c6aa64361a93f2_qw35sft2_68992a21', 'check_triple_replace__076bfd2d3ed2cd7e42a937a1dd80c8a9_qw35sft2_6ca0fcbf', 'check_dual_ext__95ee9f441436911870f312520ed4f195_qw35sft2_fe2021e4', 'check_dual_replace__e470774862e9eb18ed4fb0ad458ab39f_qw35sft2_45510beb', 'check_ext_and_wordwrap__6ac7db98d9670abbcf37d963dc27bc84_qw35sft2_dbef1ed3', 'check_mean_call__4c84b586b63f7c1cf9b6df39aedcd5ed_qw35sft2_c0f39136', 'check_dual_replace__d87d9b6cb6d6a7fc9d5d36ba0c96e0de_qw35sft2_3f330cae']

def check_row_scores__dfb627af00a19217031ee997d3cc240d(result, expected, **options):
    """Check scores in a row with partial credit per cell."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total_cells = len(expected_values)
    if total_cells == 0:
        return 0.0
    correct = 0
    for i in range(min(len(actual_values), total_cells)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        try:
            if float(actual) == float(exp):
                correct += 1
        except (TypeError, ValueError):
            if str(actual).strip() == str(exp).strip():
                correct += 1
    return correct / total_cells

def check_sorted_columns__e0472552b507a6465c494b007b9b0305(result, expected, **options):
    """Check if columns are sorted correctly by verifying cell values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not values or not expected_values:
        return 0.0
    total = len(expected_values)
    matched = 0
    for (cell_ref, exp_val) in expected_values.items():
        actual = values.get(cell_ref)
        if actual is None:
            continue
        try:
            if isinstance(exp_val, (int, float)):
                if abs(float(actual) - float(exp_val)) < 0.01:
                    matched += 1
            elif str(actual).strip().lower() == str(exp_val).strip().lower():
                matched += 1
        except (TypeError, ValueError):
            if str(actual).strip().lower() == str(exp_val).strip().lower():
                matched += 1
    return matched / total if total > 0 else 0.0

def check_earliest_entry__53a0fb3c44861b74092bea203bb878cc(result, expected, **options):
    """Check if E1 has correct header and E2 has correct paper title.
    Partial credit: header 0.3, title 0.7.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_header = expected.get('header', 'Earliest Paper')
    expected_title = expected.get('title', '')
    actual_e1 = str(result.get('e1', '')).strip().lower() if result.get('e1') else ''
    if actual_e1 == expected_header.lower():
        score += 0.3
    actual_e2 = str(result.get('e2', '')).strip().lower() if result.get('e2') else ''
    exp_title = expected_title.lower()
    if actual_e2 == exp_title:
        score += 0.7
    return min(score, 1.0)

def check_descending_sort__f8df3778414dcdc67e5f887268cde89d(result, expected, **options):
    """Check if res.txt contains numbers sorted in descending order.

    Args:
        result: local file path to res.txt (from vm_file getter)
        expected: dict with 'expected_numbers' list
    Returns:
        float: 1.0 if numbers match expected descending order, 0.0 otherwise
    """
    if not result:
        return 0.0
    try:
        with open(result) as f:
            content = f.read().strip()
    except Exception:
        return 0.0
    numbers = [int(x) for x in re.findall('\\d+', content)]
    expected_numbers = expected.get('expected_numbers', [])
    if numbers == expected_numbers:
        return 1.0
    if numbers and numbers == sorted(numbers, reverse=True):
        return 0.5
    return 0.0

def check_sar_memory_report__a8cfa9828c8ad4ee330aec94e4adaf4c(result, expected, **options):
    """Check sar memory report contains expected headers and line count.

    Partial credit:
    - 0.5 for containing all expected keywords
    - 0.5 for having at least expected number of timestamped lines
    """
    if not result or result.get('error'):
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    score = 0.0
    keywords = expected.get('keywords', [])
    if keywords:
        all_found = all((kw in content for kw in keywords))
        if all_found:
            score += 0.5
    min_lines = int(expected.get('min_timestamp_lines', 0))
    if min_lines > 0:
        time_regex = '([01]\\d|2[0-3]):[0-5]\\d:([0-5]\\d|60)'
        timestamp_lines = len(re.findall(time_regex, content))
        if timestamp_lines >= min_lines:
            score += 0.5
    return score

def check_total_row__bf5b353ae5739d8c5c981c28dcf67581(result, expected, **options):
    """Check if total row was added correctly.

    Scoring:
    - 0.50: A9 contains 'Total' (case-insensitive, partial match)
    - 0.50: D9 value equals expected sum (-210)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    a9_val = result.get('A9')
    if a9_val is not None and 'total' in str(a9_val).strip().lower():
        score += 0.5
    d9_val = result.get('D9')
    expected_sum = expected.get('expected_sum', -210)
    tolerance = expected.get('tolerance', 1.0)
    if d9_val is not None:
        try:
            d9_val = float(d9_val)
            if abs(d9_val - expected_sum) <= tolerance:
                score += 0.5
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_year_values__5d33fe5f963081b2f481ac5f73c3ccbc(result, expected, **options):
    """Check if Year column has correct header and year values.
    Partial credit: header worth 0.16, each year worth 0.168.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('column_e', [])
    expected_header = expected.get('header', 'Year')
    expected_years = expected.get('years', [])
    if not actual or len(actual) < 1:
        return 0.0
    score = 0.0
    header_val = str(actual[0]).strip().lower() if actual[0] else ''
    if header_val == expected_header.lower():
        score += 0.16
    per_year = 0.168
    for (i, exp_year) in enumerate(expected_years):
        idx = i + 1
        if idx < len(actual) and actual[idx] is not None:
            actual_str = str(actual[idx]).strip().rstrip('.0')
            exp_str = str(exp_year).strip()
            if actual_str == exp_str:
                score += per_year
    return min(score, 1.0)

def check_dir_structure__f789180956e4b0547b5929aa18a0352e(result, expected, **options):
    """Check directory structure with partial credit per subdirectory."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    required_dirs = expected.get('required_dirs', [])
    if not required_dirs:
        return 0.0
    score_per_dir = 1.0 / len(required_dirs)
    total_score = 0.0
    for d in required_dirs:
        if result.get(d, False):
            total_score += score_per_dir
    return min(total_score, 1.0)

def check_budget_reservation_info__a4a15f9cf85bf343423f5b15bc460e50(result, expected, **options):
    """Check budget.com reservation page with partial credit for URL and location."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    url = result.get('url', '')
    expected_url_pattern = expected.get('url_pattern', 'reservation')
    if expected_url_pattern.lower() in url.lower():
        score += 0.5
    if result.get('has_boston_logan', False) or result.get('has_bos_code', False):
        score += 0.5
    return min(score, 1.0)

def check_heading2_alignment__131e78d644cdd4aebe878331238d0634(result, expected, **options):
    """Check if all Heading 2 paragraphs are centered. Partial credit per heading."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    headings = result.get('headings', [])
    if not headings:
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'CENTER')
    correct = 0
    for h in headings:
        if h.get('alignment') == expected_alignment:
            correct += 1
    return correct / len(headings)

def check_main_caption__909b46690f7fc746deb97149427ce266(result, expected, **options):
    """Check if window caption matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_caption = result.get('caption', '')
    expected_caption = expected.get('expected_caption', '')
    if actual_caption and expected_caption and (actual_caption.strip() == expected_caption.strip()):
        return 1.0
    return 0.0

def check_name_and_count__009b0a2429b483dc78b74a0168744c1a(result, expected, **options):
    """Check if the most frequent name and its count are correct, with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_name = expected.get('expected_name', '')
    expected_count = expected.get('expected_count', 0)
    actual_name = result.get('name')
    if actual_name and str(actual_name).strip() == str(expected_name).strip():
        score += 0.5
    actual_count = result.get('count')
    try:
        if float(actual_count) == float(expected_count):
            score += 0.5
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_pass_fail_held__d026dfbc396400238630feee90a0bfc7(result, expected, **options):
    """Check that column D has correct Pass/Fail/Held values based on marks thresholds."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not expected_values:
        return 0.0
    correct = 0
    total = 0
    for (cell_ref, exp_val) in expected_values.items():
        total += 1
        actual = values.get(cell_ref)
        if exp_val is None:
            if actual is None or actual == '' or actual == '--':
                correct += 1
        elif actual is not None and str(actual).strip() == str(exp_val).strip():
            correct += 1
    if total == 0:
        return 0.0
    return correct / total

def check_frame_extract__5e444f1b525a6375eead73fcf3fce163(result, expected, **options):
    """Check if video frames were extracted successfully.
    Partial credit:
      0.5 - At least 1 PNG frame exists in the directory
      1.0 - Frame count is within reasonable range for a 5-second clip
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    count = result.get('frame_count', 0)
    if count == 0:
        return 0.0
    min_frames = expected.get('min_frames', 10)
    max_frames = expected.get('max_frames', 300)
    score = 0.5
    if min_frames <= count <= max_frames:
        score = 1.0
    return score

def check_user_homedir__df678257c48d0ebe8f2e03968592b783(result, expected, **options):
    """Check if user exists and has correct home directory.

    Scoring:
    - 0.5: user exists
    - 0.5: home directory matches expected
    """
    score = 0.0
    if not isinstance(result, dict) or not result.get('exists'):
        return 0.0
    score += 0.5
    expected_home = expected.get('expected_home', '')
    actual_home = result.get('home', '')
    if expected_home and actual_home == expected_home:
        score += 0.5
    return min(score, 1.0)

def check_compose_cc__457e68f7ef86811b4470193bf0ca1e7e(result, expected, **options):
    """Check if CC field contains expected email in accessibility tree."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tree_content = ''
    if isinstance(result, dict):
        tree_content = result.get('tree_content', '')
    elif isinstance(result, str):
        tree_content = result
    expected_cc = expected.get('expected_cc', '').lower()
    if not expected_cc:
        return 0.0
    if expected_cc in tree_content.lower():
        return 1.0
    return 0.0

def check_row_hidden__18d397a77a413967bac86e92a1394add(result, expected, **options):
    """Check if specified rows are hidden."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    hidden_states = result.get('hidden_states', {})
    expected_rows = expected.get('expected_hidden_rows', [])
    if not expected_rows:
        return 0.0
    correct = 0
    for row in expected_rows:
        if hidden_states.get(str(row)) is True:
            correct += 1
    return correct / len(expected_rows)

def check_navigator_visible__a0fa1d9cceca2f2896bb79deb4b8d794(result, expected, **options):
    """Check that the Navigator panel is open and visible in the accessibility tree.

    The Navigator panel in LibreOffice Impress can be opened via View > Navigator
    (Shift+Ctrl+F5). When open, it appears as a panel/window with 'Navigator'
    in its title in the accessibility tree.

    Returns 1.0 if Navigator is visible, 0.0 otherwise.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    tree_str = result if isinstance(result, str) else str(result)
    if 'Navigator' in tree_str:
        return 1.0
    return 0.0

def check_user_password_home__6fbd0fb5ef1077b146e8ceece522a6b4(result, expected, **options):
    """Check user creation with password and home directory.

    Scoring:
    - 0.33: user exists
    - 0.34: home directory correct
    - 0.33: password correct
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('user_exists'):
        score += 0.33
    if result.get('home_correct'):
        score += 0.34
    if result.get('password_correct'):
        score += 0.33
    return min(score, 1.0)

def check_picture_bottom_left__77f9e57f0af70c9666efd508929f72c3(result, expected, **options):
    """Check if picture is positioned at the bottom-left of the slide."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    pictures = result.get('pictures', [])
    if not pictures:
        return 0.0
    slide_height = result.get('slide_height', 10287000)
    slide_width = result.get('slide_width', 18288000)
    pic = pictures[0]
    pic_top = pic.get('top', 0)
    pic_bottom = pic.get('bottom', 0)
    pic_left = pic.get('left', 0)
    pic_height = pic.get('height', 0)
    score = 0.0
    bottom_threshold = slide_height * 0.7
    if pic_top >= bottom_threshold:
        score += 0.5
    elif pic_bottom >= slide_height * 0.8:
        score += 0.3
    left_threshold = slide_width * 0.4
    if pic_left < left_threshold:
        score += 0.5
    elif pic_left < slide_width * 0.5:
        score += 0.3
    return min(score, 1.0)

def check_row_scores__ba349a57c83f48325e3b8e8c17891359(result, expected, **options):
    """Check scores in a row with partial credit per cell."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total_cells = len(expected_values)
    if total_cells == 0:
        return 0.0
    correct = 0
    for i in range(min(len(actual_values), total_cells)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        try:
            if float(actual) == float(exp):
                correct += 1
        except (TypeError, ValueError):
            if str(actual).strip() == str(exp).strip():
                correct += 1
    return correct / total_cells

def check_background_color__0fc932ce1ad9456b6d290262a9551736(result, expected, **options):
    """Check if background is filled with expected color and object layer is unchanged.

    result: dict with 'result_path' and 'reference_path' (local temp files)
    expected: dict with 'color' key, value is 'red' or 'blue'

    Background pixels are identified as white/near-white pixels (all channels > 240)
    in the reference image (original has white background).

    Scoring:
      0.5 - background pixels have correct color
      0.5 - object pixels are unchanged
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    result_path = result.get('result_path', '')
    reference_path = result.get('reference_path', '')
    if not result_path or not reference_path:
        return 0.0
    try:
        result_img = Image.open(result_path).convert('RGB')
        ref_img = Image.open(reference_path).convert('RGB')
        result_pixels = np.array(result_img)
        ref_pixels = np.array(ref_img)
        (height, width) = ref_pixels.shape[:2]
        expected_color = expected.get('color', 'red')
        bg_mask = np.all(ref_pixels > 240, axis=2)
        obj_mask = ~bg_mask
        bg_count = int(np.sum(bg_mask))
        obj_count = int(np.sum(obj_mask))
        score = 0.0
        if bg_count > 0:
            bg_result = result_pixels[bg_mask]
            r = bg_result[:, 0].astype(int)
            g = bg_result[:, 1].astype(int)
            b = bg_result[:, 2].astype(int)
            if expected_color == 'red':
                correct = (r > 200) & (g < 50) & (b < 50)
            elif expected_color == 'blue':
                correct = (b > 200) & (r < 50) & (g < 50)
            else:
                correct = np.zeros(bg_count, dtype=bool)
            ratio = float(np.sum(correct)) / bg_count
            if ratio > 0.95:
                score += 0.5
        if obj_count > 0:
            obj_ref = ref_pixels[obj_mask].astype(int)
            obj_res = result_pixels[obj_mask].astype(int)
            diff = np.abs(obj_ref - obj_res)
            matching = np.all(diff <= 5, axis=1)
            ratio = float(np.sum(matching)) / obj_count
            if ratio > 0.95:
                score += 0.5
        return score
    except Exception:
        return 0.0
    finally:
        if os.path.exists(result_path):
            os.unlink(result_path)
        if os.path.exists(reference_path):
            os.unlink(reference_path)

def check_dir_move__a7ceb843bee94531ff589f6c5602590a(result, expected, **options):
    """Check directory move with partial credit.
    0.5 for dir4 existing in target (dir1), 0.5 for dir4 removed from source (dir3).
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('dir4_in_dir1', '').strip() == 'YES':
        score += 0.5
    if result.get('dir4_in_dir3', '').strip() == 'NO':
        score += 0.5
    return score

def check_column_deleted__51f83653b8f37d64288aaccfa8443c8f(result, expected, **options):
    """Check if a column was properly deleted from the spreadsheet.
    Partial credit: 0.5 for header gone, 0.5 for remaining headers correct."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    headers = result.get('headers', [])
    deleted_header = expected.get('deleted_header')
    remaining_headers = expected.get('remaining_headers')
    score = 0.0
    while headers and headers[-1] is None:
        headers.pop()
    header_strs = [str(h).strip() if h is not None else None for h in headers]
    if deleted_header not in header_strs:
        score += 0.5
    if remaining_headers:
        non_none_headers = [h for h in header_strs if h is not None]
        expected_list = remaining_headers
        if non_none_headers == expected_list:
            score += 0.5
    return min(score, 1.0)

def check_cs_data__6d4ad198dfe8ff0251a084115fc1444c(result, expected, **options):
    """Check if CS data values match expected with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_applied = expected.get('expected_applied')
    actual_applied = result.get('value_a')
    if actual_applied is not None and expected_applied is not None:
        try:
            if int(float(actual_applied)) == int(expected_applied):
                score += 0.4
        except (ValueError, TypeError):
            pass
    expected_supported = expected.get('expected_supported')
    actual_supported = result.get('value_b')
    if actual_supported is not None and expected_supported is not None:
        try:
            if int(float(actual_supported)) == int(expected_supported):
                score += 0.4
        except (ValueError, TypeError):
            pass
    header_a = result.get('header_a', '')
    header_b = result.get('header_b', '')
    if header_a and header_b:
        a_ok = any((kw in header_a.lower() for kw in ['appli', 'total', 'number', '#']))
        b_ok = any((kw in header_b.lower() for kw in ['support', 'grant', 'approved', '#']))
        if a_ok and b_ok:
            score += 0.2
        elif a_ok or b_ok:
            score += 0.1
    return min(score, 1.0)

def check_chapter_list__a4d1b434aaedff8ce21d5e5d82eac67e(result, expected, **options):
    """Check if chapter list contains expected chapters with partial credit."""
    if result.get('error'):
        return 0.0
    expected_count = expected.get('expected_count', 8)
    expected_keywords = expected.get('expected_keywords', [])
    lines = result.get('lines', [])
    score = 0.0
    if len(lines) == expected_count:
        score += 0.3
    elif abs(len(lines) - expected_count) <= 1:
        score += 0.15
    if expected_keywords:
        found = 0
        content_lower = result.get('content', '').lower()
        for kw in expected_keywords:
            if kw.lower() in content_lower:
                found += 1
        keyword_ratio = found / len(expected_keywords)
        score += 0.7 * keyword_ratio
    return min(score, 1.0)

def check_manifest_structure__783369d0e7e20c29051ab5abdc75ebb3(result, expected, **options):
    """Check manifest.json has correct name, version, description, and required feature keys.

    Partial credit:
    - 0.2 for correct name
    - 0.2 for correct version
    - 0.2 for correct description
    - 0.2 per required feature key (split evenly)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    manifest = result.get('manifest', {})
    if not isinstance(manifest, dict):
        return 0.0
    score = 0.0
    expected_name = expected.get('name', '')
    if manifest.get('name') == expected_name:
        score += 0.2
    expected_version = expected.get('version', '')
    if manifest.get('version') == expected_version:
        score += 0.2
    expected_desc = expected.get('description', '')
    if expected_desc:
        if manifest.get('description') == expected_desc:
            score += 0.2
    elif not manifest.get('description'):
        score += 0.2
    required_keys = expected.get('required_keys', [])
    if required_keys:
        key_score = 0.4 / len(required_keys)
        for key in required_keys:
            if key in manifest:
                score += key_score
    return min(score, 1.0)

def check_orientation__060b6ae42d4d041f46057b0dff06cb53(result, expected, **options):
    """Check if slide orientation matches expected."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    expected_orientation = expected.get('expected_orientation', 'portrait')
    actual_orientation = result.get('orientation', '')
    if actual_orientation == expected_orientation:
        return 1.0
    return 0.0

def check_station_removed__dc2cb3918367317ad49987158bb0de7d(result, expected, **options):
    """Check that lines with a specific station have been removed from the document."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_station', True):
        score += 0.6
    expected_remaining = expected.get('expected_remaining_lines', 90)
    actual_remaining = result.get('total_lines', 0)
    if actual_remaining == expected_remaining:
        score += 0.4
    elif abs(actual_remaining - expected_remaining) <= 2:
        score += 0.2
    return min(score, 1.0)

def check_venv__3936cc1f2257a01f73d410787af3a332(result, expected, **options):
    """Check virtual environment creation with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('venv_dir', False):
        score += 0.4
    if result.get('python_binary', False):
        score += 0.3
    if result.get('pip_binary', False):
        score += 0.3
    return min(score, 1.0)

def check_multi_replacement__43d92268ad8323e8d3f0af9cc0d7b2ce(result, expected, **options):
    """Check that multiple find-and-replace operations were performed correctly.
    Partial credit: each replacement pair is worth equal fraction.

    Expected rules:
        replacements: list of {"old_word": str, "new_word": str, "expected_count": int}
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else ''
    if not content:
        return 0.0
    replacements = expected.get('replacements', [])
    if not replacements:
        return 0.0
    score = 0.0
    weight = 1.0 / len(replacements)
    for rep in replacements:
        old_word = rep.get('old_word', '')
        new_word = rep.get('new_word', '')
        expected_count = rep.get('expected_count', 1)
        old_gone = old_word not in content if old_word else True
        new_present = content.count(new_word) >= expected_count
        if old_gone and new_present:
            score += weight
    return min(score, 1.0)

def check_optimized_bubblesort__8aa2e8e2be1ca5d610e977b9c78060b2(result, expected, **options):
    """Check sorted output correctness AND code optimization.

    Partial credit:
        0.5 - output contains correctly sorted numbers
        0.5 - code contains early-exit optimization (swapped flag)

    Args:
        result: dict with 'output' (res.txt content) and 'code' (bubbleSort.py content)
        expected: dict with 'expected_numbers' and 'optimization_keywords'
    Returns:
        float: 0.0 to 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    output_content = result.get('output', '') or ''
    if output_content:
        numbers = [int(x) for x in re.findall('\\d+', output_content)]
        expected_numbers = expected.get('expected_numbers', [])
        if numbers == expected_numbers:
            score += 0.5
    code_content = result.get('code', '') or ''
    if code_content:
        keywords = expected.get('optimization_keywords', [])
        code_lower = code_content.lower()
        if any((kw in code_lower for kw in keywords)):
            score += 0.5
    return min(score, 1.0)

def check_do_not_track__956b75b2b213e77466e1066455a2b384(result, expected, **options):
    """Check if Do Not Track is enabled in Chrome settings.

    Args:
        result: Value from enable_do_not_track getter (may be bool or string).
        expected: Dict with 'expected_enabled' key (True/False).

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, str):
        actual_enabled = result.lower() == 'true'
    else:
        actual_enabled = bool(result)
    expected_enabled = expected.get('expected_enabled', True)
    if actual_enabled == expected_enabled:
        return 1.0
    return 0.0

def check_verified_invoices__445a1e82d5445803b67308e02fdc3544(result, expected, **options):
    """Check that the correct matching invoices are in the verified folder.
    Partial credit: 0.5 per correct file present.
    Penalty: -0.25 per unexpected file (clamped to 0).
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_files = result.get('files', [])
    expected_files = expected.get('expected_files', [])
    score = 0.0
    per_file = 1.0 / max(len(expected_files), 1)
    for ef in expected_files:
        if ef in actual_files:
            score += per_file
    unexpected = [f for f in actual_files if f not in expected_files and f.endswith('.pdf')]
    score -= len(unexpected) * 0.25
    return max(0.0, min(score, 1.0))

def check_enabled_experiments_contains__1656c6075ec2ba63a026b59d9c69a38f(enabled_experiments, rule):
    """
    Check if the expected experiment names are contained in the enabled experiments list.
    Uses subset/containment check instead of strict list equality.
    """
    enabled_experiments_names = [experiment.split('@')[0] for experiment in enabled_experiments]
    if rule['type'] == 'names':
        return 1.0 if all((name in enabled_experiments_names for name in rule['names'])) else 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_notification_dual_settings__0a0313fec36cad178ad94d7913e78201(result, expected, **options):
    """Check both Do Not Disturb (show-banners=false) and Lock Screen Notifications (show-in-lock-screen=false).
    Partial credit: 0.5 for each correct setting.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_banners = expected.get('expected_show_banners', 'false')
    expected_lock = expected.get('expected_show_in_lock_screen', 'false')
    if result.get('show_banners') == expected_banners:
        score += 0.5
    if result.get('show_in_lock_screen') == expected_lock:
        score += 0.5
    return min(score, 1.0)

def check_summary_rep_totals__8a3eb44c98a42fd6dcd94122bb8b7a9a(result, expected, **options):
    """Check Summary sheet has correct per-rep totals."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not result.get('sheet_found'):
        return 0.0
    score = 0.0
    expected_totals = expected.get('expected_totals', {})
    actual_totals = result.get('rep_totals', {})
    headers = result.get('headers', [])
    if any(('rep' in h or 'name' in h or 'sales' in h for h in headers)):
        score += 0.1
    if any(('total' in h for h in headers)):
        score += 0.1
    num_reps = len(expected_totals)
    correct = 0
    for (rep_name, exp_total) in expected_totals.items():
        actual_val = actual_totals.get(rep_name)
        if actual_val is None:
            for (actual_name, actual_v) in actual_totals.items():
                if rep_name.lower() in actual_name.lower() or actual_name.lower() in rep_name.lower():
                    actual_val = actual_v
                    break
        if actual_val is not None:
            try:
                if abs(float(actual_val) - float(exp_total)) < 1:
                    correct += 1
            except (ValueError, TypeError):
                pass
    if num_reps > 0:
        score += 0.8 * (correct / num_reps)
    return min(score, 1.0)

def check_period_rate_col__02a40d6cc593a8648c91f6def88b5342(result, expected, **options):
    """Check period rate column: header + values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    header = result.get('header', '')
    expected_header = rules.get('expected_header', 'Period Rate (%)')
    if header and expected_header.lower() in str(header).lower():
        score += 0.2
    values = result.get('values', {})
    spot_checks = rules.get('spot_checks', {})
    if spot_checks:
        correct = 0
        total = len(spot_checks)
        tolerance = rules.get('tolerance', 0.01)
        for (cell, expected_val) in spot_checks.items():
            actual = values.get(cell)
            if actual is not None and abs(float(actual) - float(expected_val)) < tolerance:
                correct += 1
        if total > 0:
            score += 0.8 * (correct / total)
    return min(score, 1.0)

def check_blue_color__be59b6debe1fbee8fc53696bdb2273b8(result, expected, **options):
    """Check if BLUE color constant is correctly defined."""
    if isinstance(result, dict) and result.get('error') and (not result.get('defined')):
        return 0.0
    expected_rgb = expected.get('expected_rgb', [0, 0, 255])
    if result.get('defined') and result.get('blue') == expected_rgb:
        return 1.0
    return 0.0

def check_enabled_experiments_subset__197c5f21c248c0028a67e57d9193addd(enabled_experiments, rule, **options):
    """Check that required experiments are a subset of enabled experiments.

    Unlike the official check_enabled_experiments which uses exact list equality,
    this metric uses subset checking to verify that the required experiments
    are present in the enabled list, regardless of other pre-existing experiments.
    """
    if not isinstance(enabled_experiments, list):
        return 0.0
    enabled_experiments_names = [exp.split('@')[0] for exp in enabled_experiments]
    if rule['type'] == 'names':
        return 1.0 if set(rule['names']).issubset(set(enabled_experiments_names)) else 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_cumulative_col__6bc8ae788974152807b9723497a26da7(result, expected, **options):
    """Check cumulative sales column header and values."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header', '')
    if header and 'cumul' in str(header).lower() and ('sale' in str(header).lower()):
        score += 0.2
    expected_values = expected.get('expected_values', [])
    actual_values = result.get('values', [])
    tolerance = expected.get('tolerance', 1.0)
    if len(actual_values) == len(expected_values):
        correct = 0
        for (actual, exp) in zip(actual_values, expected_values):
            if actual is not None and exp is not None:
                try:
                    if abs(float(actual) - float(exp)) < tolerance:
                        correct += 1
                except (ValueError, TypeError):
                    pass
        score += 0.8 * (correct / len(expected_values))
    return min(score, 1.0)

def check_contrast_increase__a61095a992c03ac5e17b282d64e6f2d0(result, expected, **options):
    """Check that contrast increased while maintaining structural similarity."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    min_similarity = expected.get('min_similarity', 0.3)
    contrast_increased = result.get('contrast_increased', False)
    similarity = result.get('similarity', 0.0)
    score = 0.0
    if contrast_increased:
        score += 0.5
    if similarity >= min_similarity:
        score += 0.5
    return min(score, 1.0)

def check_top_performers__85190e1339a2cb1782e74f01270bc172(result, expected, **options):
    """Check TopPerformers sheet has correct month-top rep mapping."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not result.get('sheet_found'):
        return 0.0
    score = 0.0
    expected_map = expected.get('expected_top', {})
    actual_map = result.get('month_top', {})
    if result.get('sheet_found'):
        score += 0.1
    if result.get('row_count', 0) >= len(expected_map):
        score += 0.1
    num_months = len(expected_map)
    correct = 0
    for (month, exp_rep) in expected_map.items():
        actual_rep = actual_map.get(month)
        if actual_rep is None:
            for (am, ar) in actual_map.items():
                if am.lower() == month.lower():
                    actual_rep = ar
                    break
        if actual_rep and exp_rep:
            exp_parts = exp_rep.lower().split()
            actual_parts = actual_rep.lower().split()
            if exp_rep.lower() == actual_rep.lower():
                correct += 1
            elif any((p in actual_parts for p in exp_parts)):
                correct += 0.8
    if num_months > 0:
        score += 0.8 * (correct / num_months)
    return min(score, 1.0)

def check_invoice_total__aa615bdde918e746d9793d851a2bb4f8(result, expected, **options):
    """Check that the total.txt file contains the correct total invoice amount.
    Accepts various formats: 12160, 12160.00, $12,160.00, etc.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    expected_total = expected.get('expected_total', 0)
    cleaned = content.replace(',', '').replace('$', '').replace(' ', '')
    numbers = re.findall('[\\d]+(?:\\.[\\d]+)?', cleaned)
    for num_str in numbers:
        try:
            value = float(num_str)
            if abs(value - expected_total) < 1.0:
                return 1.0
        except ValueError:
            continue
    return 0.0

def check_do_not_track__8d6140655fd3ffdbb0f2973a3475732a(result, expected, **options):
    """Check if Chrome's Do Not Track setting matches expected state."""
    expected_value = expected.get('enabled', True)
    if isinstance(result, str):
        result_bool = result.lower() == 'true'
    elif isinstance(result, bool):
        result_bool = result
    else:
        return 0.0
    return 1.0 if result_bool == expected_value else 0.0

def check_do_not_track__baa6ff2cd46140d8d50ba44da41af676(result, expected, **options):
    """Check if Do Not Track is enabled.
    The enable_do_not_track getter returns a string ('true'/'false'),
    so we normalize both sides before comparison.
    """
    expected_value = expected.get('expected_value', True)
    if isinstance(result, str):
        result_bool = result.lower() == 'true'
    elif isinstance(result, bool):
        result_bool = result
    else:
        return 0.0
    if isinstance(expected_value, str):
        expected_bool = expected_value.lower() == 'true'
    elif isinstance(expected_value, bool):
        expected_bool = expected_value
    else:
        return 0.0
    return 1.0 if result_bool == expected_bool else 0.0

def check_sum_row__99d61f333e953d709c5fe5f221bfb63e(result, expected, **options):
    """Check if sum row has correct label and values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    label = result.get('label')
    if label is not None and str(label).strip().lower() == 'total':
        score += 0.2
    expected_sales = expected.get('expected_sales')
    actual_sales = result.get('sales_total')
    if actual_sales is not None and expected_sales is not None:
        try:
            if abs(float(actual_sales) - float(expected_sales)) < 1.0:
                score += 0.4
        except (TypeError, ValueError):
            pass
    expected_cogs = expected.get('expected_cogs')
    actual_cogs = result.get('cogs_total')
    if actual_cogs is not None and expected_cogs is not None:
        try:
            if abs(float(actual_cogs) - float(expected_cogs)) < 1.0:
                score += 0.4
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_compose_recipients__9a5e3321e7daf20a2f35cb35407bc19e(result, expected, **options):
    """Check compose recipients with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tree_content = ''
    if isinstance(result, dict):
        tree_content = result.get('tree_content', '')
    elif isinstance(result, str):
        tree_content = result
    tree_lower = tree_content.lower()
    score = 0.0
    expected_to = expected.get('expected_to', '').lower()
    if expected_to and expected_to in tree_lower:
        score += 0.5
    expected_cc = expected.get('expected_cc', '').lower()
    if expected_cc and expected_cc in tree_lower:
        score += 0.5
    return min(score, 1.0)

def check_sales_by_rep__5c7590b2fce2859cb43c07916cdde25c(result, expected, **options):
    """Check that Sheet2 contains correct total sales per sales rep."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    sales_by_rep = result.get('sales_by_rep', {})
    expected_sales = expected.get('expected_sales', {})
    if not sales_by_rep or not expected_sales:
        return 0.0
    score = 0.0
    total_reps = len(expected_sales)
    per_rep_score = 1.0 / total_reps if total_reps > 0 else 0.0
    for (rep, exp_val) in expected_sales.items():
        actual_val = sales_by_rep.get(rep)
        if actual_val is not None:
            if abs(float(actual_val) - float(exp_val)) < 0.01:
                score += per_rep_score
    return min(score, 1.0)

def check_profit_column__11a8fd5475e2aba6ae95fda21befab56(result, expected, **options):
    """Check if Profit column has correct header and computed values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header')
    if header is not None and str(header).strip().lower() == 'profit':
        score += 0.2
    expected_profits = expected.get('expected_profits', [])
    actual_values = result.get('values', [])
    if len(actual_values) >= 10 and len(expected_profits) >= 10:
        correct_count = 0
        for i in range(10):
            actual = actual_values[i]
            exp = expected_profits[i]
            if actual is not None and exp is not None:
                try:
                    if abs(float(actual) - float(exp)) < 1.0:
                        correct_count += 1
                except (TypeError, ValueError):
                    pass
        score += 0.8 * (correct_count / 10.0)
    return min(score, 1.0)

def check_rename_copy__25257c2d579c2a114cd376d84d920386(result, expected, **options):
    """Check rename + copy operations with partial credit.

    Expected output lines from getter command:
    FILE1_GONE:YES/NO
    DATA_IN_DIR1:YES/NO
    DATA_IN_DIR2:YES/NO
    DATA_IN_DIR3:YES/NO
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    output = result if isinstance(result, str) else str(result)
    score = 0.0
    if 'FILE1_GONE:YES' in output:
        score += 0.25
    if 'DATA_IN_DIR1:YES' in output:
        score += 0.25
    if 'DATA_IN_DIR2:YES' in output:
        score += 0.25
    if 'DATA_IN_DIR3:YES' in output:
        score += 0.25
    return min(score, 1.0)

def check_do_not_track__7317fcab347872285f8a0f137be4010c(result, expected, **options):
    """Check if Do Not Track setting matches expected value."""
    expected_value = expected.get('enabled', True)
    if isinstance(result, str):
        result_bool = result.lower() == 'true'
    else:
        result_bool = bool(result)
    return 1.0 if result_bool == expected_value else 0.0

def check_novel_info__6a9459e3699067e60bf00ce031fcbd18(result, expected, **options):
    """Check if novel info file contains correct title and chapter count. Partial credit."""
    if not isinstance(result, dict) or result.get('error') or (not result.get('exists')):
        return 0.0
    content = result.get('content', '')
    content_lower = result.get('content_lower', '')
    score = 0.0
    expected_title = expected.get('expected_title', 'pass through').lower()
    if expected_title in content_lower:
        score += 0.5
    expected_count = str(expected.get('expected_count', 5))
    if expected_count in content:
        score += 0.5
    return min(score, 1.0)

def check_net_income_column__bef74e7adbdaaf2c559cbe5b3bf80940(result, expected, **options):
    """Check if column C contains correct Net Income values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_header = expected.get('expected_header')
    actual_header = result.get('header')
    if actual_header is not None and expected_header is not None:
        if str(actual_header).strip().lower() == str(expected_header).strip().lower():
            score += 0.2
    expected_values = expected.get('expected_values', [])
    actual_values = result.get('values', [])
    if expected_values and actual_values:
        matches = 0
        total = len(expected_values)
        for i in range(min(len(actual_values), total)):
            actual = actual_values[i]
            exp = expected_values[i]
            if actual is not None and exp is not None:
                try:
                    if abs(float(actual) - float(exp)) < 0.01:
                        matches += 1
                except (ValueError, TypeError):
                    pass
        if total > 0:
            score += 0.8 * (matches / total)
    return min(score, 1.0)

def check_contact_export__091acaa52f197bd2f2171a56782c0841(result, expected, **options):
    """Check CSV contains new contact and original contacts with partial credit."""
    import csv
    if result is None:
        return 0.0
    score = 0.0
    try:
        with open(result) as f:
            reader = csv.DictReader(f)
            rows = list(reader)
        min_rows = expected.get('min_rows', 30)
        if len(rows) >= min_rows:
            score += 0.5
        target_first = expected.get('first_name', '')
        target_last = expected.get('last_name', '')
        target_email = expected.get('email', '')
        for row in rows:
            first_match = row.get('First Name', '').strip() == target_first
            last_match = row.get('Last Name', '').strip() == target_last
            email_match = target_email in row.get('Primary Email', '')
            if first_match and last_match and email_match:
                score += 0.5
                break
    except Exception:
        return 0.0
    return min(score, 1.0)

def check_row_values__4551fb3f098ec9fe882ac83568cd431f(result, expected, **options):
    """Check if the row values match expected values with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('label_found'):
        return 0.0
    expected_values = expected.get('expected_values', {})
    tolerance = expected.get('tolerance', 0.5)
    if not expected_values:
        return 0.0
    total_checks = len(expected_values)
    passed = 0
    for (col, exp_val) in expected_values.items():
        actual = result.get('values', {}).get(col)
        if actual is None:
            continue
        try:
            if abs(float(actual) - float(exp_val)) <= tolerance:
                passed += 1
        except (ValueError, TypeError):
            continue
    return passed / total_checks if total_checks > 0 else 0.0

def check_columns_hidden__a4f4a609a17cec62166f8b2d532de58f(result, expected, **options):
    """Check if specified columns are hidden. Partial credit per column."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    hidden_status = result.get('hidden_status', {})
    expected_hidden = expected.get('expected_hidden', {})
    if not expected_hidden:
        return 0.0
    correct = 0
    total = len(expected_hidden)
    for (col, should_be_hidden) in expected_hidden.items():
        if hidden_status.get(col) == should_be_hidden:
            correct += 1
    return correct / total if total > 0 else 0.0

def check_row_values__0cd3deb0e30597bee67bd1ee61842fcf(result, expected, **options):
    """Check if the row values match expected values with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('label_found'):
        return 0.0
    expected_values = expected.get('expected_values', {})
    tolerance = expected.get('tolerance', 0.5)
    if not expected_values:
        return 0.0
    total_checks = len(expected_values)
    passed = 0
    for (col, exp_val) in expected_values.items():
        actual = result.get('values', {}).get(col)
        if actual is None:
            continue
        try:
            if abs(float(actual) - float(exp_val)) <= tolerance:
                passed += 1
        except (ValueError, TypeError):
            continue
    return passed / total_checks if total_checks > 0 else 0.0

def check_column_sorted_ascending__0f45568bd733d9527d872ad7c053c2c8(result, expected, **options):
    """Check if text column values are sorted in ascending alphabetical order."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', [])
    if not values:
        return 0.0
    str_values = [v for v in values if v is not None]
    if len(str_values) != len(values):
        return 0.0
    expected_sorted = sorted(str_values, key=lambda x: x.lower() if isinstance(x, str) else str(x))
    if str_values == expected_sorted:
        return 1.0
    correct = sum((1 for (a, b) in zip(str_values, expected_sorted) if a == b))
    return correct / len(expected_sorted) * 0.5

def check_import_exists__d8b2f0ccd335f9a3a3203363805b5ac0(result, expected, **options):
    """Check if the expected import statement exists in the file content."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    content = result.get('content', '')
    expected_import = expected.get('expected_import', '')
    if expected_import and expected_import in content:
        return 1.0
    return 0.0

def check_badge_deleted__6b4bf280e2d30f49e851d7a1a643ef02(result, expected, **options):
    """Check if the badge image has been deleted from slide 1."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    badge_exists = result.get('badge_exists', True)
    expected_exists = expected.get('badge_should_exist', False)
    if badge_exists == expected_exists:
        return 1.0
    return 0.0

def check_merged_novel__2abe2760620ddfe22faad3f384e1cae4(result, expected, **options):
    """Check if merged novel file exists and contains all chapters. Partial credit per chapter."""
    if not isinstance(result, dict) or result.get('error') or (not result.get('exists')):
        return 0.0
    score = 0.0
    if result.get('has_ch1'):
        score += 0.2
    if result.get('has_ch2'):
        score += 0.2
    if result.get('has_ch3'):
        score += 0.2
    if result.get('has_ch4'):
        score += 0.2
    if result.get('has_ch5'):
        score += 0.2
    min_lines = expected.get('min_lines', 2000)
    if result.get('line_count', 0) < min_lines:
        score *= 0.5
    return min(score, 1.0)

def check_dir_structure__95c704f58705a26de50af292b2eb5059(result, expected, **options):
    """Check directory structure exists. Partial credit per directory."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_subdirs = expected.get('expected_subdirs', [])
    total_checks = 1 + len(expected_subdirs)
    score = 0.0
    if result.get('base_exists'):
        score += 1.0 / total_checks
    subdirs = result.get('subdirs', {})
    for sd in expected_subdirs:
        if subdirs.get(sd):
            score += 1.0 / total_checks
    return min(score, 1.0)

def check_model_names__55e51bced75ba3e5e08e1be190df564e(result, expected, **options):
    """Check if the text file contains all expected model names with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else str(result)
    content_lower = content.lower()
    expected_names = expected.get('expected_names', [])
    if not expected_names:
        return 0.0
    found = 0
    for name in expected_names:
        if name.lower() in content_lower:
            found += 1
    return found / len(expected_names)

def check_window_size__e47348d3e93e8a085dab2f1753ad88d5(result, expected, **options):
    """Check if WIDTH and HEIGHT match expected values. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_width = expected.get('expected_width', 800)
    expected_height = expected.get('expected_height', 600)
    if result.get('width') == expected_width:
        score += 0.5
    if result.get('height') == expected_height:
        score += 0.5
    return score

def check_song_list_content__09209b5893b5a25bd669841ecee9c784(result, expected, **options):
    """Check if song_list.txt contains all expected MP3 filenames. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    score = 0.0
    per_file = 1.0 / len(expected_files)
    for f in expected_files:
        if f in content:
            score += per_file
    return min(round(score, 2), 1.0)

def check_zip_contents__b6394280513f3b0b8d3eaa6ea7ea0992(result, expected, **options):
    """Check if zip contains expected files. Partial credit: 0.5 for zip existing, 0.5 for correct contents."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if not result.get('exists', False):
        return 0.0
    score += 0.3
    expected_files = set(expected.get('expected_files', []))
    actual_files = set(result.get('files', []))
    if not expected_files:
        return score
    matches = len(expected_files & actual_files)
    score += 0.7 * (matches / len(expected_files))
    return min(score, 1.0)

def check_compose_subject__df6702fb4aa059fb4a86e567f6f9f62d(result, expected, **options):
    """Check if Thunderbird compose window title contains expected subject."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    window_title = ''
    if isinstance(result, dict):
        window_title = result.get('window_title', '')
    elif isinstance(result, str):
        window_title = result
    expected_subject = expected.get('expected_subject', '')
    if not expected_subject:
        return 0.0
    if expected_subject.lower() in window_title.lower():
        return 1.0
    return 0.0

def check_grayscale__71dd5cb242259598df523724839bf008(result, expected, **options):
    """Check if the image is grayscale (desaturated).
    Verifies that R ≈ G ≈ B for the vast majority of pixels.
    """
    if result is None:
        return 0.0
    try:
        img = Image.open(result)
        if img.mode == 'L':
            return 1.0
        img_rgb = img.convert('RGB')
        pixels = np.array(img_rgb)
        r = pixels[:, :, 0].astype(int)
        g = pixels[:, :, 1].astype(int)
        b = pixels[:, :, 2].astype(int)
        diff_rg = np.abs(r - g)
        diff_rb = np.abs(r - b)
        diff_gb = np.abs(g - b)
        max_diff = np.maximum(np.maximum(diff_rg, diff_rb), diff_gb)
        tolerance = expected.get('tolerance', 10)
        grayscale_ratio = np.mean(max_diff <= tolerance)
        logging.debug(f'Grayscale ratio: {grayscale_ratio:.4f} (tolerance={tolerance})')
        if grayscale_ratio >= 0.95:
            return 1.0
        elif grayscale_ratio >= 0.8:
            return 0.5
        else:
            return 0.0
    except Exception as e:
        logging.error(f'Error checking grayscale: {e}')
        return 0.0

def check_dir_ownership__425163a454302c61dc411d05ac9420b7(result, expected, **options):
    """Check directory ownership and permissions.

    Scoring:
    - 0.34: owner matches
    - 0.33: group matches
    - 0.33: permissions match
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_owner = expected.get('expected_owner', '')
    expected_group = expected.get('expected_group', '')
    expected_perms = expected.get('expected_permissions', '')
    if expected_owner and result.get('owner') == expected_owner:
        score += 0.34
    if expected_group and result.get('group') == expected_group:
        score += 0.33
    if expected_perms and result.get('permissions') == expected_perms:
        score += 0.33
    return min(score, 1.0)

def check_do_not_track__fef997178e549d845fc31b37602ff4d0(result, expected, **options):
    """Check if Do Not Track setting matches expected value.

    Scoring:
    - 1.0 if Do Not Track enabled state matches expected
    - 0.0 otherwise
    """
    expected_enabled = expected.get('expected_enabled', True)
    if isinstance(result, bool):
        return 1.0 if result == expected_enabled else 0.0
    if isinstance(result, str):
        actual = result.lower() == 'true'
        return 1.0 if actual == expected_enabled else 0.0
    return 0.0

def check_descending_order__40424b53798c2f42a45c9b5d1f78c9a0(result, expected, **options):
    """Check if column values are sorted in descending order and first value matches."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', [])
    if not values:
        return 0.0
    score = 0.0
    expected_first = expected.get('expected_first_value')
    try:
        if expected_first is not None and abs(float(values[0]) - float(expected_first)) < 0.01:
            score += 0.5
    except (TypeError, ValueError):
        pass
    numeric_vals = []
    for v in values:
        if v is not None:
            try:
                numeric_vals.append(float(v))
            except (TypeError, ValueError):
                pass
    if len(numeric_vals) >= 2:
        is_descending = all((numeric_vals[i] >= numeric_vals[i + 1] for i in range(len(numeric_vals) - 1)))
        if is_descending:
            score += 0.5
    return min(score, 1.0)

def check_delimiter_replaced__0a313c15fc4fb10c48c753cea47e4527(result, expected, **options):
    """Check that delimiter replacement was performed correctly."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    old_count = result.get('old_delimiter_count', -1)
    if old_count == 0:
        score += 0.5
    expected_new_count = expected.get('expected_new_delimiter_count', 300)
    new_count = result.get('new_delimiter_count', 0)
    if new_count == expected_new_count:
        score += 0.3
    elif abs(new_count - expected_new_count) <= 5:
        score += 0.15
    parts_count = result.get('parts_with_new_delim', 0)
    if parts_count == 4:
        score += 0.2
    return min(score, 1.0)

def check_sorted_names__56c311ddff224f797b6614d5151a3c35(result, expected, **options):
    """Check if employee names are sorted alphabetically A-Z."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    names = result.get('names', [])
    expected_names = expected.get('sorted_names', [])
    if not names or not expected_names:
        return 0.0
    matches = 0
    total = len(expected_names)
    for i in range(min(len(names), total)):
        if names[i] == expected_names[i]:
            matches += 1
    return matches / total if total > 0 else 0.0

def check_mpstat_report__28ae802c5923a6e7ecfebd0e5cc59702(result, expected, **options):
    """Check mpstat per-core CPU report contains expected headers and line count.

    Partial credit:
    - 0.5 for containing all expected keywords
    - 0.5 for having at least expected number of timestamped lines
    """
    if not result or result.get('error'):
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    score = 0.0
    keywords = expected.get('keywords', [])
    if keywords:
        all_found = all((kw in content for kw in keywords))
        if all_found:
            score += 0.5
    min_lines = int(expected.get('min_timestamp_lines', 0))
    if min_lines > 0:
        time_regex = '([01]\\d|2[0-3]):[0-5]\\d:([0-5]\\d|60)'
        timestamp_lines = len(re.findall(time_regex, content))
        if timestamp_lines >= min_lines:
            score += 0.5
    return score

def check_sorted_values__0c918bbe7b80a9e03be9fde4824c7824(result, expected, **options):
    """Check if column values are sorted alphabetically. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', [])
    expected_values = expected.get('expected_sorted_values', [])
    if not values or not expected_values:
        return 0.0
    correct = 0
    total = max(len(expected_values), len(values))
    for i in range(min(len(values), len(expected_values))):
        if values[i].strip() == expected_values[i].strip():
            correct += 1
    return correct / total if total > 0 else 0.0

def check_title_style_bg__d816da51e2714a5b86252b84b6e77125(result, expected, **options):
    """Check title bold/italic and background color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_bold') is True:
        score += 0.33
    if result.get('title_italic') is True:
        score += 0.33
    expected_bg = expected.get('bg_color', '0000FF')
    actual_bg = result.get('bg_color', '')
    if actual_bg and expected_bg.upper() in actual_bg.upper():
        score += 0.34
    return min(score, 1.0)

def check_title_added__2130fa798c3c78ff755136f88d8b3efd(result, expected, **options):
    """Check that a title line was added as the first line of the document."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_title = expected.get('expected_title', 'Train Records')
    first_line = result.get('first_line', '')
    if expected_title.lower() in first_line.lower():
        score += 0.6
    expected_total = expected.get('expected_total_lines', 101)
    actual_total = result.get('total_lines', 0)
    if actual_total >= expected_total:
        score += 0.4
    return min(score, 1.0)

def check_fps_value__6d1f5614df32aa00fba65a5a2ca43cb2(result, expected, **options):
    """Check if FPS value matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_fps = expected.get('expected_fps', 20)
    if result.get('fps') == expected_fps:
        return 1.0
    return 0.0

def check_title_centered__f845858287fa4f02431e180c3025ad53(result, expected, **options):
    """Check if the title paragraph is centered."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if result.get('is_centered') is True:
        return 1.0
    return 0.0

def check_strikethrough__a345b8271c8198e8c95ec913d78c781f(result, expected, **options):
    """Check if target text has strikethrough formatting applied."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    target_texts = expected.get('target_texts', [])
    if not target_texts:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(target_texts)
    for target in target_texts:
        for para in paragraphs:
            if target.lower() in para['text'].lower():
                if para['strikethrough']:
                    score += per_item
                break
    return min(score, 1.0)

def check_do_not_track_enabled__2a127fa9e4cd79198d970c24eb95cab0(result, expected, **options):
    """Check if Do Not Track setting matches expected value."""
    if result is None:
        return 0.0
    expected_value = expected.get('expected_value', 'true')
    result_str = str(result).lower()
    expected_str = str(expected_value).lower()
    if result_str == expected_str:
        return 1.0
    return 0.0

def check_word_counts__4ca76e68b8033d8d7fa4f889ee4aceef(result, expected, **options):
    """Check that column values match expected word counts.
    Uses partial credit: each correct row contributes equally.
    Handles both int and float (e.g., 3 vs 3.0)."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None:
                try:
                    if int(float(actual)) == int(exp_val):
                        correct += 1
                except (ValueError, TypeError):
                    pass
    return correct / total if total > 0 else 0.0

def check_booking_form_state__cf1c94ef66b2087660e9bbd9fc7c9d7e(result, expected, **options):
    """Check booking page URL and form inputs with partial credit.

    Scoring:
    - 0.4: On the correct booking page (URL contains expected pattern)
    - 0.3: Name field matches expected
    - 0.3: Email field matches expected
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    actual_url = result.get('url', '')
    expected_url_pattern = expected.get('url_pattern', '')
    if expected_url_pattern and expected_url_pattern in actual_url:
        score += 0.4
    actual_name = str(result.get('name', '')).strip()
    expected_name = str(expected.get('expected_name', '')).strip()
    if expected_name and actual_name.lower() == expected_name.lower():
        score += 0.3
    actual_mail = str(result.get('mail', '')).strip()
    expected_mail = str(expected.get('expected_mail', '')).strip()
    if expected_mail and actual_mail.lower() == expected_mail.lower():
        score += 0.3
    return min(score, 1.0)

def check_screen_lock_settings__39d04f914d2dbd8627ddcd7f791e785b(result, expected, **options):
    """Check screen lock settings with partial credit.

    Gives 0.5 for correct idle delay, 0.5 for correct lock-enabled.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_idle = expected.get('expected_idle_delay', '')
    expected_lock = expected.get('expected_lock_enabled', '')
    if expected_idle and expected_idle in result.get('idle_delay', ''):
        score += 0.5
    if result.get('lock_enabled', '').strip() == expected_lock:
        score += 0.5
    return min(score, 1.0)

def check_zone3_row_totals__c68f755a470f57d0018a142ca4ed4f38(result, expected, **options):
    """Check if Zone 3 product row totals are correct. Partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    cells = ['F17', 'F18', 'F19']
    per_cell = 1.0 / len(cells)
    for cell in cells:
        actual = result.get(cell)
        exp = expected.get(cell)
        if actual is not None and exp is not None:
            try:
                if abs(float(actual) - float(exp)) < 0.01:
                    score += per_cell
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_chart_exists__9affe8a59575f48a48ea3570040f007c(result, expected, **options):
    """Check if at least one chart of the expected type exists with correct data references."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    chart_count = result.get('chart_count', 0)
    expected_min = expected.get('min_charts', 1)
    expected_type = expected.get('chart_type', None)
    expected_value_col = expected.get('expected_value_column', None)
    expected_cat_col = expected.get('expected_category_column', None)
    score = 0.0
    if chart_count >= expected_min:
        score += 0.34
    elif chart_count > 0:
        score += 0.17
    type_matched_idx = None
    if expected_type and result.get('chart_types'):
        for (idx, ct) in enumerate(result['chart_types']):
            if expected_type.lower() in ct.lower():
                type_matched_idx = idx
                score += 0.33
                break
    chart_series_info = result.get('chart_series_info', [])
    if expected_value_col or expected_cat_col:
        check_idx = type_matched_idx if type_matched_idx is not None else 0
        if check_idx < len(chart_series_info):
            series_list = chart_series_info[check_idx]
            data_score = 0.0
            value_ok = False
            cat_ok = False
            for s in series_list:
                val_ref = s.get('value_ref', '')
                if expected_value_col and val_ref:
                    col_pattern = '\\$' + re.escape(expected_value_col.upper()) + '\\$'
                    if re.search(col_pattern, val_ref):
                        value_ok = True
                cat_ref = s.get('category_ref', '')
                if expected_cat_col and cat_ref:
                    col_pattern = '\\$' + re.escape(expected_cat_col.upper()) + '\\$'
                    if re.search(col_pattern, cat_ref):
                        cat_ok = True
            if expected_value_col and expected_cat_col:
                if value_ok and cat_ok:
                    data_score = 0.33
                elif value_ok or cat_ok:
                    data_score = 0.165
            elif expected_value_col:
                data_score = 0.33 if value_ok else 0.0
            elif expected_cat_col:
                data_score = 0.33 if cat_ok else 0.0
            score += data_score
    return min(score, 1.0)

def check_column_sorted_descending__3201def6075c63c69d0100c10729d61d(result, expected, **options):
    """Check if column values are sorted in descending order."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', [])
    if not values:
        return 0.0
    numeric_values = [v for v in values if v is not None]
    if len(numeric_values) != len(values):
        return 0.0
    expected_sorted = sorted(numeric_values, reverse=True)
    if numeric_values == expected_sorted:
        return 1.0
    correct = sum((1 for (a, b) in zip(numeric_values, expected_sorted) if a == b))
    return correct / len(expected_sorted) * 0.5

def check_range_values__160ec07bd1e641be9b31cf92f3848b19(result, expected, **options):
    """Check if multiple cell values match expected values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (cell, exp_val) in expected_values.items():
        actual = actual_values.get(cell)
        if actual is not None and str(actual).strip() == str(exp_val).strip():
            correct += 1
    return correct / total if total > 0 else 0.0

def check_grocery_entry__00c370b29bbdeefb624743c61a277c9e(result, expected, **options):
    """Check if grocery expense entry was added correctly.

    Scoring:
    - 0.50: Amount in D9 is approximately -186.93
    - 0.25: Type in C9 is 'Expense' (case-insensitive)
    - 0.25: Description in A9 is not empty
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_amount = expected.get('expected_amount', -186.93)
    tolerance = expected.get('tolerance', 0.5)
    actual_amount = result.get('D')
    if actual_amount is not None:
        try:
            actual_amount = float(actual_amount)
            if abs(actual_amount - expected_amount) <= tolerance:
                score += 0.5
        except (ValueError, TypeError):
            pass
    actual_type = result.get('C')
    if actual_type is not None and str(actual_type).strip().lower() == 'expense':
        score += 0.25
    actual_desc = result.get('A')
    if actual_desc is not None and str(actual_desc).strip():
        score += 0.25
    return min(score, 1.0)

def check_title_bg__c2ac2f94e6c641c5a5c2e733ff74ade1(result, expected, **options):
    """Check title text and background color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title_text', '')
    if expected_title and expected_title.lower() == actual_title.lower():
        score += 0.5
    expected_colors = expected.get('expected_bg_colors', [])
    actual_color = result.get('bg_color', '')
    if actual_color and expected_colors:
        actual_upper = actual_color.upper()
        for ec in expected_colors:
            if ec.upper() == actual_upper:
                score += 0.5
                break
    return min(score, 1.0)

def check_model_count__204b2db58c306983ce47277dbc39788e(result, expected, **options):
    """Check if the text file contains the expected model count."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '').strip() if isinstance(result, dict) else str(result).strip()
    expected_count = str(expected.get('expected_count', ''))
    try:
        actual_num = int(content.strip())
        expected_num = int(expected_count)
        if actual_num == expected_num:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_empty_count__4128b5119649f756a2c7d9fc4321d25f(result, expected, **options):
    """Check if result.txt contains the correct count of empty stock prices."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    content = result.get('content', '').strip()
    if not content:
        return 0.0
    expected_count = str(expected.get('expected_count', ''))
    if expected_count and expected_count in content:
        return 1.0
    return 0.0

def check_total_row__6a822179da68f351e1e0a8c5b6636775(result, expected, **options):
    """Check total row has correct label and sums."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    label = result.get('label', '')
    if label and 'total' in str(label).lower():
        score += 0.2
    expected_sales = expected.get('expected_total_sales')
    actual_sales = result.get('total_sales')
    tolerance = expected.get('tolerance', 1.0)
    if actual_sales is not None and expected_sales is not None:
        try:
            if abs(float(actual_sales) - float(expected_sales)) < tolerance:
                score += 0.4
        except (ValueError, TypeError):
            pass
    expected_cogs = expected.get('expected_total_cogs')
    actual_cogs = result.get('total_cogs')
    if actual_cogs is not None and expected_cogs is not None:
        try:
            if abs(float(actual_cogs) - float(expected_cogs)) < tolerance:
                score += 0.4
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_conference_continents__a9d0754cd398d25ba1c2a300dbecfcf0(result, expected, **options):
    """Check conference continents with partial credit per row.
    Expected format: {"continents": ["North America", "Asia", ...]}
    """
    if result.get('error'):
        return 0.0
    values = result.get('values', [])
    expected_continents = expected.get('continents', [])
    if not values or not expected_continents:
        return 0.0
    n = min(len(values), len(expected_continents))
    if n == 0:
        return 0.0
    correct = 0
    for i in range(n):
        actual = values[i]
        if actual is None:
            continue
        actual_lower = actual.lower().strip()
        exp_lower = expected_continents[i].lower().strip()
        if actual_lower == exp_lower:
            correct += 1
    return correct / len(expected_continents)

def check_manifest_structure__ab85f9bbe05040b98eb1a964cbee7866(result, expected, **options):
    """Check manifest.json has correct name, version, and required feature keys.

    Partial credit:
    - 0.2 for correct name
    - 0.2 for correct version
    - 0.6 split evenly across required feature keys
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    manifest = result.get('manifest', {})
    if not isinstance(manifest, dict):
        return 0.0
    score = 0.0
    expected_name = expected.get('name', '')
    if manifest.get('name') == expected_name:
        score += 0.2
    expected_version = expected.get('version', '')
    if manifest.get('version') == expected_version:
        score += 0.2
    required_keys = expected.get('required_keys', [])
    if required_keys:
        key_score = 0.6 / len(required_keys)
        for key in required_keys:
            if key in manifest:
                score += key_score
    return min(score, 1.0)

def check_title_format__110ae328b957703cc4b8db2a2b04df5c(result, expected, **options):
    """Check if title formatting matches expected bold/italic state.
    Partial credit: 0.5 for correct bold, 0.5 for correct italic.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_bold = expected.get('expected_bold', False)
    expected_italic = expected.get('expected_italic', False)
    actual_bold = result.get('bold', False)
    actual_italic = result.get('italic', False)
    if actual_bold == expected_bold:
        score += 0.5
    if actual_italic == expected_italic:
        score += 0.5
    return score

def check_manifest_structure__61a5e99298254bc3ea7017f23352ee68(result, expected, **options):
    """Check manifest.json has correct name, version, and required feature keys.

    Partial credit:
    - 0.25 for correct name
    - 0.25 for correct version
    - 0.25 per required feature key (split evenly)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    manifest = result.get('manifest', {})
    if not isinstance(manifest, dict):
        return 0.0
    score = 0.0
    expected_name = expected.get('name', '')
    if manifest.get('name') == expected_name:
        score += 0.25
    expected_version = expected.get('version', '')
    if manifest.get('version') == expected_version:
        score += 0.25
    required_keys = expected.get('required_keys', [])
    if required_keys:
        key_score = 0.5 / len(required_keys)
        for key in required_keys:
            if key in manifest:
                score += key_score
    return min(score, 1.0)

def check_row_values__8c5594571c596e8d15ba168f23066d41(result, expected, **options):
    """Check row values with partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_vals = expected.get('expected_values', {})
    label_cell = expected.get('label_cell')
    label_value = expected.get('label_value')
    if not expected_vals:
        return 0.0
    total_checks = len(expected_vals)
    if label_cell and label_value:
        total_checks += 1
    score = 0.0
    if label_cell and label_value:
        actual_label = result.get(label_cell)
        if actual_label is not None and str(actual_label).strip().lower() == label_value.strip().lower():
            score += 1.0
    tolerance = expected.get('tolerance', 0.01)
    for (cell, exp_val) in expected_vals.items():
        actual = result.get(cell)
        if actual is None:
            continue
        try:
            actual_num = float(actual)
            exp_num = float(exp_val)
            if abs(actual_num - exp_num) <= abs(exp_num * tolerance) + 0.001:
                score += 1.0
        except (ValueError, TypeError):
            continue
    return score / total_checks if total_checks > 0 else 0.0

def check_repair_entry__c6556bc153411ec2613fae626891d4a6(result, expected, **options):
    """Check if repair expense entry was added correctly.

    Scoring:
    - 0.50: Amount in D9 is approximately -154.06
    - 0.25: Type in C9 is 'Expense' (case-insensitive)
    - 0.25: Description in A9 is not empty
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_amount = expected.get('expected_amount', -154.06)
    tolerance = expected.get('tolerance', 0.5)
    actual_amount = result.get('D')
    if actual_amount is not None:
        try:
            actual_amount = float(actual_amount)
            if abs(actual_amount - expected_amount) <= tolerance:
                score += 0.5
        except (ValueError, TypeError):
            pass
    actual_type = result.get('C')
    if actual_type is not None and str(actual_type).strip().lower() == 'expense':
        score += 0.25
    actual_desc = result.get('A')
    if actual_desc is not None and str(actual_desc).strip():
        score += 0.25
    return min(score, 1.0)

def check_column_values__6d4300c0b784865a64e946f955774003(result, expected, **options):
    """Check if column values match expected values with partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    correct = 0
    total = len(expected_values)
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None and exp_val is not None:
                try:
                    if abs(float(actual) - float(exp_val)) < 0.01:
                        correct += 1
                except (ValueError, TypeError):
                    if str(actual) == str(exp_val):
                        correct += 1
    return correct / total if total > 0 else 0.0

def check_heading_alignment__e69f89d4207a469ba2d4f747df580e7b(result, expected, **options):
    """Check if heading alignment matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('alignment', '')
    expected_align = expected.get('expected_alignment', 'center')
    if actual.lower() == expected_align.lower():
        return 1.0
    return 0.0

def check_units_by_product__8ec7487c1752ceb36e17eb91a13203df(result, expected, **options):
    """Check that Sheet2 contains correct total units per product."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    units_by_product = result.get('units_by_product', {})
    expected_units = expected.get('expected_units', {})
    if not units_by_product or not expected_units:
        return 0.0
    score = 0.0
    total_products = len(expected_units)
    per_product_score = 1.0 / total_products if total_products > 0 else 0.0
    for (product, exp_val) in expected_units.items():
        actual_val = units_by_product.get(product)
        if actual_val is not None:
            if abs(float(actual_val) - float(exp_val)) < 0.01:
                score += per_product_score
    return min(score, 1.0)

def check_uncomment__cd3236440d08e0a8882e8811427d01d3(result, expected, **options):
    """Check if the print statement has been uncommented."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    is_commented = result.get('is_commented', True)
    if not is_commented:
        return 1.0
    return 0.0

def check_vlookup_prices__40371a0bf41fabc605f9425ce4121622(result, expected, **options):
    """Check VLOOKUP results: header + price values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    total_weight = 1.0
    header = result.get('header')
    expected_header = expected.get('expected_header', 'Retail Price')
    if header and isinstance(header, str) and (expected_header.lower() in header.lower()):
        score += 0.2
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    tolerance = expected.get('tolerance', 0.01)
    if not expected_values or not actual_values:
        return score
    correct = 0
    total = len(expected_values)
    for (i, exp_val) in enumerate(expected_values):
        if i >= len(actual_values):
            break
        act_val = actual_values[i]
        if act_val is None or exp_val is None:
            if act_val == exp_val:
                correct += 1
            continue
        try:
            if abs(float(act_val) - float(exp_val)) <= tolerance:
                correct += 1
        except (TypeError, ValueError):
            continue
    if total > 0:
        score += 0.8 * (correct / total)
    return min(score, total_weight)

def check_upper_titles__6d702aff56732f3b1cb6a80e7f740e70(result, expected, **options):
    """Check that column values match expected uppercase cleaned titles.
    Uses partial credit: each correct row contributes equally."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None and str(actual).strip() == str(exp_val).strip():
                correct += 1
    return correct / total if total > 0 else 0.0

def check_scaling_factor__d5d03d50dfc17bfe4f646f3d3c9160fb(result, expected, **options):
    """Check if the text scaling factor matches the expected value with float tolerance."""
    if not result or not isinstance(result, str):
        return 0.0
    try:
        actual = float(result.strip())
    except (ValueError, TypeError):
        return 0.0
    expected_factor = expected.get('expected_factor', 1.25)
    if abs(actual - float(expected_factor)) < 0.01:
        return 1.0
    return 0.0

def check_tally_entry__89c74c0ac9693815f604c967708f8f82(result, expected, **options):
    """Check that a specific tally entry exists with correct values.

    Expected keys: service, month_contains, amount
    Partial credit: 0.33 for service match, 0.33 for month match, 0.34 for amount match.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    rows = result.get('rows', [])
    exp_service = str(expected.get('service', '')).strip().upper()
    exp_month_contains = str(expected.get('month_contains', ''))
    exp_amount = float(expected.get('amount', 0))
    best_score = 0.0
    for row in rows:
        score = 0.0
        service = str(row.get('service', '')).strip().upper()
        month = str(row.get('month', ''))
        amount = row.get('amount')
        if service == exp_service:
            score += 0.33
        if exp_month_contains and exp_month_contains in month:
            score += 0.33
        if amount is not None:
            try:
                if abs(float(amount) - exp_amount) < 0.01:
                    score += 0.34
            except (ValueError, TypeError):
                pass
        best_score = max(best_score, score)
    return min(best_score, 1.0)

def check_power_settings_combined__e637b95c625a88b3bb2cee1c77022b86(result, expected, **options):
    """Check both idle-dim disabled and screen blank set to never. Partial credit: 0.5 each."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    idle_dim = result.get('idle_dim', '')
    expected_dim = expected.get('expected_idle_dim', 'false')
    if idle_dim == expected_dim:
        score += 0.5
    idle_delay = result.get('idle_delay', '')
    expected_delay = expected.get('expected_idle_delay', 'uint32 0')
    if idle_delay == expected_delay:
        score += 0.5
    return min(score, 1.0)

def check_title_alignment__4f408a61496e9bff25fa5074f0f2e380(result, expected, **options):
    """Check if the title paragraph has the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'center')
    actual_alignment = result.get('alignment', '')
    if actual_alignment.lower() == expected_alignment.lower():
        return 1.0
    return 0.0

def check_import_extraction__c191ddd519b4043306a4effe99ae0aac(result, expected, **options):
    """Check that the file contains valid, unique, sorted import statements."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    if not lines:
        return 0.0
    score = 0.0
    if len(lines) >= 1:
        score += 0.3
    import_lines = result.get('import_lines', [])
    non_import_lines = result.get('non_import_lines', [])
    if len(non_import_lines) == 0 and len(import_lines) > 0:
        score += 0.3
    elif len(import_lines) > len(lines) * 0.8:
        score += 0.15
    if not result.get('has_duplicates', True):
        score += 0.2
    if result.get('is_sorted', False):
        score += 0.2
    return min(score, 1.0)

def check_column_values__4065e6a4b9dbccc0390b8c5a38e2d1b7(result, expected, **options):
    """Check if column cell values match expected values. Supports partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    cells = result.get('cells', {})
    expected_cells = expected.get('expected_cells', {})
    if not expected_cells:
        return 0.0
    total = len(expected_cells)
    matched = 0
    for (cell_ref, expected_val) in expected_cells.items():
        actual_val = cells.get(cell_ref)
        if actual_val is not None and expected_val is not None:
            if actual_val.strip() == str(expected_val).strip():
                matched += 1
    return matched / total if total > 0 else 0.0

def check_settings_speed_blocksize__4b546ac4b7d86d470dd227a5c08930e6(result, expected, **options):
    """Check if game speed and block size match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('game_speed') == expected.get('expected_speed'):
        score += 0.5
    if result.get('block_size') == expected.get('expected_blocksize'):
        score += 0.5
    return score

def check_sorted_sales__38a9aeefb36e58a6a0dcece70b9af3ed(result, expected, **options):
    """Check if data is sorted by Sales in descending order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    sales_values = result.get('sales_values', [])
    expected_order = expected.get('expected_sales_order', [])
    if not sales_values or len(sales_values) < 10:
        return 0.0
    is_descending = True
    for i in range(len(sales_values) - 1):
        if sales_values[i] is None or sales_values[i + 1] is None:
            is_descending = False
            break
        if float(sales_values[i]) < float(sales_values[i + 1]):
            is_descending = False
            break
    if is_descending:
        score += 0.5
    if len(expected_order) >= 10:
        correct_count = 0
        for i in range(10):
            actual = sales_values[i]
            exp = expected_order[i]
            if actual is not None and exp is not None:
                try:
                    if abs(float(actual) - float(exp)) < 1.0:
                        correct_count += 1
                except (TypeError, ValueError):
                    pass
        score += 0.5 * (correct_count / 10.0)
    return min(score, 1.0)

def check_model_list_unordered__68e1e5b5d5b054d1e9ca98b707844de8(result, rules) -> float:
    """
    Compare a parsed URL modelList against expected models, ignoring order.
    result: dict from active_tab_url_parse, e.g. {"modelList": ["iphone-16-pro-max", ...]}
    rules: dict with "expected_models" key containing list of expected model slugs.
    Returns 1.0 if all expected models are present (as a set), 0.0 otherwise.
    """
    logger.info(f'[DEBUG] check_model_list_unordered called with result: {result}')
    logger.info(f'[DEBUG] check_model_list_unordered called with rules: {rules}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    expected_models = rules.get('expected_models', [])
    actual_models = result.get('modelList', [])
    if not expected_models:
        logger.info('[DEBUG] No expected_models defined, returning 0.0')
        return 0.0
    if not actual_models:
        logger.info('[DEBUG] No modelList in result, returning 0.0')
        return 0.0
    expected_set = set(expected_models)
    actual_set = set(actual_models)
    logger.info(f'[DEBUG] Expected set: {expected_set}')
    logger.info(f'[DEBUG] Actual set: {actual_set}')
    if expected_set == actual_set:
        logger.info('[DEBUG] Sets match, returning 1.0')
        return 1.0
    else:
        logger.info(f'[DEBUG] Sets do not match. Missing: {expected_set - actual_set}, Extra: {actual_set - expected_set}')
        return 0.0

def check_half_rate_col__fd58351f5fafe47d305d3595c5202dcd(result, expected, **options):
    """Check half rate column: header + values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    header = result.get('header', '')
    expected_header = rules.get('expected_header', 'Half Rate')
    if header and expected_header.lower() in str(header).lower():
        score += 0.2
    values = result.get('values', {})
    spot_checks = rules.get('spot_checks', {})
    if spot_checks:
        correct = 0
        total = len(spot_checks)
        tolerance = rules.get('tolerance', 0.01)
        for (cell, expected_val) in spot_checks.items():
            actual = values.get(cell)
            if actual is not None and abs(float(actual) - float(expected_val)) < tolerance:
                correct += 1
        if total > 0:
            score += 0.8 * (correct / total)
    return min(score, 1.0)

def check_rows_hidden__cb276526eec2ec94231e7dedd20b5069(result, expected, **options):
    """Check if specified rows are hidden. Partial credit per row."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    hidden_states = result.get('hidden_states', {})
    expected_rows = expected.get('expected_hidden_rows', [])
    if not expected_rows:
        return 0.0
    correct = 0
    for row in expected_rows:
        if hidden_states.get(str(row)) is True:
            correct += 1
    return correct / len(expected_rows)

def check_fv_column__5c1e2dfe8bffbdb20d88aa7499677c46(result, expected, **options):
    """Check Future Value column header and computed values with tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    expected_header = rules.get('header', 'Future Value')
    actual_header = result.get('F1')
    if actual_header and str(actual_header).strip().lower() == expected_header.lower():
        score += 0.2
    expected_values = rules.get('values', {})
    tolerance = rules.get('tolerance', 0.5)
    for (cell_ref, exp_val) in expected_values.items():
        actual = result.get(cell_ref)
        if actual is not None:
            try:
                actual_num = float(actual)
                exp_num = float(exp_val)
                if exp_num != 0 and abs(actual_num - exp_num) / abs(exp_num) < tolerance:
                    score += 0.2
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_utils_function__597b85b971b86604637b8eea77967812(result, expected, **options):
    """Check utils.py with partial credit.

    Scoring:
      0.5 - File exists and contains 'def add' function definition
      0.5 - Function correctly returns sum (add(2, 3) == 5)
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('file_exists') and result.get('has_add_function'):
        score += 0.5
    expected_output = expected.get('expected_output', '5')
    actual_output = str(result.get('func_output', '')).strip()
    if actual_output == expected_output:
        score += 0.5
    return min(score, 1.0)

def check_column_sums__427c1b1ebceb5c2ead27edb218ee7bf9(result, expected, **options):
    """Check if column values match expected sums with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total = len(expected_values)
    if total == 0:
        return 0.0
    correct = 0
    for i in range(min(len(actual_values), total)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        try:
            if abs(float(actual) - float(exp)) < 0.01:
                correct += 1
        except (TypeError, ValueError):
            continue
    return correct / total

def check_year_values__dcb80f39266c26be10a70749c6e4e4f7(result, expected, **options):
    """Check if years 2021, 2020, 2019 were entered in B3, B4, B5 with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_years = expected.get('expected_years', {})
    for (cell_key, expected_year) in expected_years.items():
        actual = result.get(cell_key)
        if actual is not None:
            try:
                if int(float(actual)) == int(expected_year):
                    score += 1.0 / len(expected_years)
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_do_not_track_enabled__66408ac42f837225d56e2756a129da0a(result, expected, **options):
    """Check if Do Not Track is enabled. The getter returns a string 'true'/'false'."""
    expected_value = expected.get('enabled', 'true')
    if isinstance(result, str) and isinstance(expected_value, str):
        return 1.0 if result.lower() == expected_value.lower() else 0.0
    return 1.0 if str(result).lower() == str(expected_value).lower() else 0.0

def check_net_income_column__bc1bb048c0eeb30a1fa5f9604d4f0a11(result, expected, **options):
    """Check Net Income column header and values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header', '')
    if header and 'net income' in str(header).lower().replace('_', ' '):
        score += 0.2
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if len(actual_values) == len(expected_values) and len(expected_values) > 0:
        correct = 0
        for (a, e) in zip(actual_values, expected_values):
            try:
                if abs(float(a) - float(e)) < 0.01:
                    correct += 1
            except (TypeError, ValueError):
                pass
        score += 0.8 * (correct / len(expected_values))
    return min(score, 1.0)

def check_code_stats__7636602600807a890f06c684e5901e16(result, expected, **options):
    """Check that the stats file has exactly 3 lines of positive integers."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    parsed = result.get('parsed_numbers', [])
    if not lines:
        return 0.0
    score = 0.0
    if len(lines) == 3:
        score += 0.3
    elif len(lines) >= 1:
        score += 0.1
    valid_numbers = [n for n in parsed if n is not None and n > 0]
    if len(valid_numbers) == len(lines) and len(lines) == 3:
        score += 0.3
    elif len(valid_numbers) > 0:
        score += 0.1
    if len(parsed) >= 1 and parsed[0] is not None:
        if 5 <= parsed[0] <= 100:
            score += 0.2
    if len(parsed) >= 3 and parsed[1] is not None and (parsed[2] is not None):
        if parsed[1] > parsed[2] > 0:
            score += 0.2
    return min(score, 1.0)

def check_names_and_city__95be0722e22941aff21eff3d680c0c01(result, expected, **options):
    """Check restaurant names in A2:A6 and city values in E column. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_names = expected.get('expected_names', [])
    expected_city = expected.get('expected_city', '').lower().strip()
    expected_header = expected.get('expected_e1_header', '').lower().strip()
    actual_names = result.get('names', [])
    for (i, exp_name) in enumerate(expected_names):
        if i < len(actual_names) and actual_names[i] is not None:
            actual_lower = actual_names[i].lower().strip()
            exp_lower = exp_name.lower().strip()
            if exp_lower in actual_lower or actual_lower in exp_lower:
                score += 0.1
    actual_header = (result.get('e1_header') or '').lower().strip()
    if actual_header and expected_header in actual_header:
        score += 0.1
    actual_cities = result.get('city_values', [])
    for i in range(5):
        if i < len(actual_cities) and actual_cities[i] is not None:
            if expected_city in actual_cities[i].lower().strip():
                score += 0.08
    return min(score, 1.0)

def check_move_symlinks__2a402a2f74167209d0ebc1db72b76fd8(result, expected, **options):
    """Check move + symlink operations with partial credit.

    Expected output lines from getter command:
    FILE_IN_DIR1:YES/NO
    FILE1_REMOVED:YES/NO
    SYMLINK_DIR2:YES/NO
    SYMLINK_DIR3:YES/NO
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    output = result if isinstance(result, str) else str(result)
    score = 0.0
    if 'FILE_IN_DIR1:YES' in output:
        score += 0.25
    if 'FILE1_REMOVED:YES' in output:
        score += 0.25
    if 'SYMLINK_DIR2:YES' in output:
        score += 0.25
    if 'SYMLINK_DIR3:YES' in output:
        score += 0.25
    return min(score, 1.0)

def check_lower_titles__56ea06753821989747521678ff3594b3(result, expected, **options):
    """Check that column values match expected lowercase cleaned titles.
    Uses partial credit: each correct row contributes equally."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None and str(actual).strip() == str(exp_val).strip():
                correct += 1
    return correct / total if total > 0 else 0.0

def check_row_scores__1700c01d5ad41ab8401de0bafb801990(result, expected, **options):
    """Check scores in a row with partial credit per cell."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total_cells = len(expected_values)
    if total_cells == 0:
        return 0.0
    correct = 0
    for i in range(min(len(actual_values), total_cells)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        try:
            if float(actual) == float(exp):
                correct += 1
        except (TypeError, ValueError):
            if str(actual).strip() == str(exp).strip():
                correct += 1
    return correct / total_cells

def check_column_sorted__34bc68338cfc2efa5ed7f80b3fe4afc3(result, expected, **options):
    """Check if column values are sorted in the specified order.
    Partial credit: 0.5 for having all values, 0.5 for correct order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', [])
    order = expected.get('order', 'descending')
    expected_count = expected.get('expected_count')
    score = 0.0
    if expected_count is not None and len(values) == expected_count:
        score += 0.5
    elif expected_count is None and len(values) > 0:
        score += 0.5
    if len(values) >= 2:
        numeric_vals = []
        for v in values:
            try:
                numeric_vals.append(float(v))
            except (ValueError, TypeError):
                return score
        if order == 'descending':
            is_sorted = all((numeric_vals[i] >= numeric_vals[i + 1] for i in range(len(numeric_vals) - 1)))
        else:
            is_sorted = all((numeric_vals[i] <= numeric_vals[i + 1] for i in range(len(numeric_vals) - 1)))
        if is_sorted:
            score += 0.5
    return min(score, 1.0)

def check_title_color__6b28d17bfe07bb684917211f85bf497c(result, expected, **options):
    """Check if the title font color matches the expected color."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_color = expected.get('expected_color', '000000')
    colors = result.get('colors', [])
    if not colors:
        return 0.0
    matching = sum((1 for c in colors if c and c.upper() == expected_color.upper()))
    if matching == len(colors):
        return 1.0
    elif matching > 0:
        return matching / len(colors)
    return 0.0

def check_inbox_state__8e4a0a40e4d045486197a43a70e9b321(result, expected, **options):
    """Check inbox and trash state after email deletion."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    inbox_count = result.get('inbox_message_count', -1)
    trash_count = result.get('trash_message_count', -1)
    expected_inbox = expected.get('expected_inbox_count', 1)
    expected_trash_empty = expected.get('expected_trash_empty', True)
    if inbox_count == expected_inbox:
        score += 0.6
    if expected_trash_empty:
        if trash_count == 0:
            score += 0.4
    else:
        score += 0.4
    return min(score, 1.0)

def check_zone2_col_totals__dfa00147f689f96c837b3a22754c016b(result, expected, **options):
    """Check if Zone 2 quarterly column totals are correct. Partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    cells = ['B13', 'C13', 'D13', 'E13']
    per_cell = 1.0 / len(cells)
    for cell in cells:
        actual = result.get(cell)
        exp = expected.get(cell)
        if actual is not None and exp is not None:
            try:
                if abs(float(actual) - float(exp)) < 0.01:
                    score += per_cell
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_column_values__2b83d943849f7c327e47bb940fe75967(result, expected, **options):
    """Check if column values match expected values with partial credit per cell."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    n = len(expected_values)
    if len(actual_values) != n:
        return 0.0
    score = 0.0
    per_cell = 1.0 / n
    for (actual, exp) in zip(actual_values, expected_values):
        if actual is None:
            continue
        try:
            if abs(float(actual) - float(exp)) < 0.5:
                score += per_cell
        except (TypeError, ValueError):
            continue
    return min(score, 1.0)

def check_ascending_order__5f28938c5d9137ac59ddda20f23ca5f2(result, expected, **options):
    """Check if column values are sorted in ascending order."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', [])
    if not values:
        return 0.0
    numeric_values = []
    for v in values:
        if v is not None:
            try:
                numeric_values.append(float(v))
            except (ValueError, TypeError):
                return 0.0
    if len(numeric_values) != len(expected.get('expected_values', [])):
        pass
    expected_sorted = expected.get('expected_values', [])
    if len(numeric_values) == len(expected_sorted):
        matches = sum((1 for (a, e) in zip(numeric_values, expected_sorted) if abs(float(a) - float(e)) < 0.01))
        return matches / len(expected_sorted)
    is_ascending = all((numeric_values[i] <= numeric_values[i + 1] for i in range(len(numeric_values) - 1)))
    return 1.0 if is_ascending else 0.0

def check_bibtex_entry__dd385e64b1e5f06bf4b4c074020cb6a0(result, expected, **options):
    """Check if references.bib contains a valid BibTeX entry for the expected paper.
    Partial credit:
    - 0.5 for having a valid BibTeX entry (@article, @inproceedings, etc.)
    - 0.5 for the entry containing the expected title keywords
    """
    if not result or result.get('error'):
        return 0.0
    content = result.get('content', '').strip()
    if not content:
        return 0.0
    score = 0.0
    bibtex_pattern = '@\\w+\\s*\\{[^}]*'
    if re.search(bibtex_pattern, content):
        score += 0.5
    title_keywords = expected.get('title_keywords', [])
    if title_keywords:
        content_lower = content.lower()
        matched = sum((1 for kw in title_keywords if kw.lower() in content_lower))
        if len(title_keywords) > 0 and matched / len(title_keywords) >= 0.6:
            score += 0.5
    return min(score, 1.0)

def check_conference_countries__70ea2799be65bb426ef2e5f3f76ece43(result, expected, **options):
    """Check conference countries with partial credit per row.
    Expected format: {"countries": [["USA", "United States"], ...]}
    Each entry is a list of acceptable alternatives.
    """
    if result.get('error'):
        return 0.0
    values = result.get('values', [])
    expected_countries = expected.get('countries', [])
    if not values or not expected_countries:
        return 0.0
    n = min(len(values), len(expected_countries))
    if n == 0:
        return 0.0
    correct = 0
    for i in range(n):
        actual = values[i]
        if actual is None:
            continue
        actual_lower = actual.lower().strip()
        alternatives = expected_countries[i]
        if isinstance(alternatives, str):
            alternatives = [alternatives]
        for alt in alternatives:
            if alt.lower().strip() == actual_lower:
                correct += 1
                break
    return correct / len(expected_countries)

def check_column_values__312ae9a7a4fdef961589d00f7ec76217(result, expected, **options):
    """Check if column values match expected values with partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    correct = 0
    total = len(expected_values)
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None and exp_val is not None:
                try:
                    if abs(float(actual) - float(exp_val)) < 0.01:
                        correct += 1
                except (ValueError, TypeError):
                    if str(actual) == str(exp_val):
                        correct += 1
    return correct / total if total > 0 else 0.0

def check_row_values__b46e00d9bcf11154ef195dbbf421e737(result, expected, **options):
    """Check row values with partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_vals = expected.get('expected_values', {})
    label_cell = expected.get('label_cell')
    label_value = expected.get('label_value')
    if not expected_vals:
        return 0.0
    total_checks = len(expected_vals)
    if label_cell and label_value:
        total_checks += 1
    score = 0.0
    if label_cell and label_value:
        actual_label = result.get(label_cell)
        if actual_label is not None and str(actual_label).strip().lower() == label_value.strip().lower():
            score += 1.0
    tolerance = expected.get('tolerance', 0.01)
    for (cell, exp_val) in expected_vals.items():
        actual = result.get(cell)
        if actual is None:
            continue
        try:
            actual_num = float(actual)
            exp_num = float(exp_val)
            if abs(actual_num - exp_num) <= abs(exp_num * tolerance) + 0.001:
                score += 1.0
        except (ValueError, TypeError):
            continue
    return score / total_checks if total_checks > 0 else 0.0

def check_function_rename__bef206759cc52e3fdba71a4f0e77e8b3(result, expected, **options):
    """Check if function was renamed correctly with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else str(result)
    expected_func = expected.get('expected_function', '')
    old_func = expected.get('old_function', '')
    import re
    score = 0.0
    if f'def {expected_func}' in content:
        score += 0.5
    if old_func and f'def {old_func}' not in content:
        score += 0.25
    call_pattern = re.compile(f'(?<!def\\s){re.escape(expected_func)}\\(\\)')
    if call_pattern.search(content):
        score += 0.25
    return min(score, 1.0)

def check_horizontal_flip__30cadb0ec7eca28df693cf924ab88301(result, expected, **options):
    """
    Metric that checks if the result image is a horizontal flip of the original.
    The getter already computed SSIM between the result and the flipped original.
    Returns 1.0 if SSIM >= threshold, 0.0 otherwise.
    """
    if result.get('error'):
        return 0.0
    ssim_score = result.get('ssim', 0.0)
    threshold = expected.get('ssim_threshold', 0.9)
    if ssim_score >= threshold:
        return 1.0
    return 0.0

def check_row_values__767f7417cb0fc080051f9d7103e674cb(result, expected, **options):
    """Check row values with partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_vals = expected.get('expected_values', {})
    label_cell = expected.get('label_cell')
    label_value = expected.get('label_value')
    if not expected_vals:
        return 0.0
    total_checks = len(expected_vals)
    if label_cell and label_value:
        total_checks += 1
    score = 0.0
    if label_cell and label_value:
        actual_label = result.get(label_cell)
        if actual_label is not None and str(actual_label).strip().lower() == label_value.strip().lower():
            score += 1.0
    tolerance = expected.get('tolerance', 0.01)
    for (cell, exp_val) in expected_vals.items():
        actual = result.get(cell)
        if actual is None:
            continue
        try:
            actual_num = float(actual)
            exp_num = float(exp_val)
            if abs(actual_num - exp_num) <= abs(exp_num * tolerance) + 0.01:
                score += 1.0
        except (ValueError, TypeError):
            continue
    return score / total_checks if total_checks > 0 else 0.0

def check_no_pictures__48004e0f145da635f1d1ef72ec0c75ce(result, expected, **options):
    """Check that there are no pictures on the slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    picture_count = result.get('picture_count', -1)
    expected_count = expected.get('expected_count', 0)
    if picture_count == expected_count:
        return 1.0
    return 0.0

def check_notes_bg__7bae80483f423bfa3b3bbb7f9e9eafdb(result, expected, **options):
    """Check notes text and background color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_notes = expected.get('expected_notes', '')
    actual_notes = result.get('notes_text', '')
    if expected_notes and expected_notes in actual_notes:
        score += 0.5
    expected_colors = expected.get('expected_bg_colors', [])
    actual_color = result.get('bg_color', '')
    if actual_color and expected_colors:
        actual_upper = actual_color.upper()
        for ec in expected_colors:
            if ec.upper() == actual_upper:
                score += 0.5
                break
    return min(score, 1.0)

def check_color_and_size__043ea3777ab786b45abdbd440171ffc3(result, expected, **options):
    """Check background color AND image size with partial credit.

    result: dict with 'result_path', 'reference_path', 'width', 'height'
    expected: dict with 'color', 'expected_width', 'expected_height'

    Scoring: 0.5 for correct background color, 0.5 for correct size.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    result_path = result.get('result_path', '')
    reference_path = result.get('reference_path', '')
    expected_width = expected.get('expected_width', 0)
    expected_height = expected.get('expected_height', 0)
    if result.get('width') == expected_width and result.get('height') == expected_height:
        score += 0.5
    try:
        result_img = Image.open(result_path)
        ref_img = Image.open(reference_path)
        result_pixels = np.array(result_img)
        ref_pixels = np.array(ref_img)
        expected_color = expected.get('color', 'blue')
        color_correct = True
        (h, w) = result_pixels.shape[:2]
        sample_positions = [(0, 0), (0, w - 1), (h - 1, 0), (h - 1, w - 1), (0, w // 2), (h - 1, w // 2), (h // 2, 0), (h // 2, w - 1)]
        for (y, x) in sample_positions:
            (r, g, b) = result_pixels[y, x][:3]
            if expected_color == 'blue':
                if not (int(b) > int(r) and int(b) > int(g)):
                    color_correct = False
                    break
            elif expected_color == 'red':
                if not (int(r) > int(g) and int(r) > int(b)):
                    color_correct = False
                    break
        if color_correct:
            score += 0.5
    except Exception:
        pass
    finally:
        if os.path.exists(result_path):
            os.unlink(result_path)
        if os.path.exists(reference_path):
            os.unlink(reference_path)
    return min(score, 1.0)

def check_dir_exists__93ffa0d39e3e07dab53be514a31da362(result, expected, **options):
    """Check if the directory exists based on getter output."""
    if not result or isinstance(result, str):
        return 0.0
    if result.get('error'):
        return 0.0
    should_exist = expected.get('should_exist', True)
    actually_exists = result.get('exists', False)
    if should_exist == actually_exists:
        return 1.0
    return 0.0

def check_gemini_response__ae954eaf7f0ec05c9c2918c4eec73bcc(result, expected, **options):
    """Check if the text file contains key phrases from the expected Gemini response with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else str(result)
    key_phrases = expected.get('key_phrases', [])
    if not key_phrases:
        return 0.0
    found = 0
    for phrase in key_phrases:
        if phrase.lower() in content.lower():
            found += 1
    return found / len(key_phrases)

def check_number_format__092734e9cc221c6f71e099e209498fb4(result, expected, **options):
    """Check if cells have the expected number format (2 decimal places)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    formats = result.get('formats', {})
    if not formats:
        return 0.0
    expected_format = expected.get('expected_format', '0.00')
    total = len(formats)
    if total == 0:
        return 0.0
    correct = 0
    for (cell_ref, fmt) in formats.items():
        if fmt == expected_format:
            correct += 1
    return correct / total

def check_currency_format__6f34d4ab7744f84c76fe24f84ac55b09(result, expected, **options):
    """Check if cells have currency format with $ symbol and 2 decimal places."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    formats = result.get('formats', [])
    if not formats:
        return 0.0
    matches = 0
    total = len(formats)
    for fmt in formats:
        fmt_str = str(fmt).upper()
        has_dollar = '$' in fmt_str
        has_two_decimals = '.00' in fmt_str or '0.00' in fmt_str
        if has_dollar and has_two_decimals:
            matches += 1
    return matches / total if total > 0 else 0.0

def check_padded_ids__598fab012ada3ac788280b75e5c01391(result, expected, **options):
    """Check that column D contains Old IDs padded to 8 digits with leading zeros."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total = len(expected_values)
    matches = 0
    for i in range(min(len(actual_values), total)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        actual_clean = str(actual).strip()
        if actual_clean == exp:
            matches += 1
    return matches / total if total > 0 else 0.0

def check_sorted_ascending__85a76f9ca092376f5f612fc019e6af39(result, expected, **options):
    """Check that Old ID column is sorted in ascending numerical order."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_sorted', [])
    if not actual_values:
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_count', 29)
    if len(actual_values) == expected_count:
        score += 0.2
    is_sorted = all((actual_values[i] <= actual_values[i + 1] for i in range(len(actual_values) - 1)))
    if is_sorted:
        score += 0.4
    if expected_values and len(actual_values) == len(expected_values):
        matches = sum((1 for (a, e) in zip(actual_values, expected_values) if a == e))
        score += 0.4 * (matches / len(expected_values))
    return min(score, 1.0)

def check_highest_price__4c5e4ba8ba69ff54d1e2443439646d44(result, expected, **options):
    """Check if result.txt contains the correct highest price info. Partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    content = result.get('content', '').lower()
    if not content:
        return 0.0
    score = 0.0
    rules = expected
    expected_symbol = rules.get('expected_symbol', '').lower()
    expected_price = str(rules.get('expected_price', ''))
    if expected_symbol and expected_symbol in content:
        score += 0.5
    if expected_price and expected_price in content:
        score += 0.5
    return min(score, 1.0)

def check_columns_filled_and_sum__f0c5ce9b9ab74cdf70d3ea75ec500ce3(result, expected, **options):
    """Check that columns B and D are filled and B31 SUM is correct.

    Partial scoring:
    - 0.33: Column B sample cell (B15) has a numeric value
    - 0.34: Column D sample cell (D15) has a numeric value
    - 0.33: B31 SUM is within tolerance of expected value
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    values = result.get('values', {})
    score = 0.0
    b15 = values.get('B15')
    if b15 is not None:
        try:
            float(b15)
            score += 0.33
        except (TypeError, ValueError):
            pass
    d15 = values.get('D15')
    if d15 is not None:
        try:
            float(d15)
            score += 0.34
        except (TypeError, ValueError):
            pass
    b31 = values.get('B31')
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    if b31 is not None and expected_value is not None:
        try:
            actual_num = float(b31)
            expected_num = float(expected_value)
            if abs(actual_num - expected_num) <= tolerance:
                score += 0.33
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_numbered_notes__5e14eceba731d7ebe21b244525a2edc1(result, expected, **options):
    """Check that notes with slide numbers match expected format.

    Expected format per line: 'Slide N: note text'
    Only slides that have notes should appear.
    """
    if result.get('error'):
        return 0.0
    actual_lines = result.get('lines', [])
    expected_entries = expected.get('expected_entries', [])
    if not expected_entries:
        return 0.0
    total = len(expected_entries)
    matched = 0
    for (i, exp_entry) in enumerate(expected_entries):
        if i < len(actual_lines):
            actual = actual_lines[i].strip().lower()
            exp = exp_entry.strip().lower()
            if actual == exp:
                matched += 1
    return matched / total if total > 0 else 0.0

def check_zone1_row_totals__b15b7a7665c71e3a9519658d5a13d57f(result, expected, **options):
    """Check if Zone 1 product row totals are correct. Partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    cells = ['F3', 'F4', 'F5']
    per_cell = 1.0 / len(cells)
    for cell in cells:
        actual = result.get(cell)
        exp = expected.get(cell)
        if actual is not None and exp is not None:
            try:
                if abs(float(actual) - float(exp)) < 0.01:
                    score += per_cell
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_backup_move__1a7983e090de51b7f1a1b54e751935d9(result, expected, **options):
    """Check backup + move operations with partial credit.

    Expected output lines from getter command:
    BACKUP_EXISTS:YES/NO
    BACKUP_CONTENT_MATCH:YES/NO
    FILE1_IN_DIR2:YES/NO
    FILE1_REMOVED:YES/NO
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    output = result if isinstance(result, str) else str(result)
    score = 0.0
    if 'BACKUP_EXISTS:YES' in output:
        score += 0.25
    if 'BACKUP_CONTENT_MATCH:YES' in output:
        score += 0.25
    if 'FILE1_IN_DIR2:YES' in output:
        score += 0.25
    if 'FILE1_REMOVED:YES' in output:
        score += 0.25
    return min(score, 1.0)

def check_sidebar_hidden__c002e9fc79af5f293a68e76c7f778a2e(result, expected, **options):
    """Check that the Properties sidebar on the right is NOT visible in the accessibility tree.

    The sidebar in LibreOffice Impress shows panels like 'Properties', 'Layouts',
    'Master Slide' etc. When closed via View > Sidebar (Ctrl+F5), these elements
    disappear from the accessibility tree.

    Returns 1.0 if sidebar is hidden, 0.0 if still visible.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    tree_str = result if isinstance(result, str) else str(result)
    sidebar_indicators = ['Properties', 'Layouts', 'Master Slide:']
    found_count = 0
    for indicator in sidebar_indicators:
        if indicator in tree_str:
            found_count += 1
    if found_count >= 2:
        return 0.0
    return 1.0

def check_vtt_conversion__9ac6e58c407c5a10e99c8a824933ab2f(result, expected, **options):
    """Check VTT conversion result with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('file_exists'):
        score += 0.34
    if result.get('has_webvtt_header'):
        score += 0.33
    if result.get('has_cues'):
        score += 0.33
    return min(score, 1.0)

def check_aws_total__d9fdc5409ef8a872907b1ddc3db9cd54(result, expected, **options):
    """Check if a total row for AWS charges has been added to the tally book.

    Looks for a numeric value matching expected_total (with tolerance) in column C
    below the original 6 data rows (header + 5 entries).
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    rows = result.get('rows', [])
    expected_total = expected.get('expected_total')
    tolerance = expected.get('tolerance', 0.05)
    if expected_total is None:
        return 0.0
    for row in rows[6:]:
        if row is None:
            continue
        for cell_val in row:
            if cell_val is None:
                continue
            try:
                val = float(cell_val)
                if abs(val - expected_total) < tolerance:
                    return 1.0
            except (TypeError, ValueError):
                continue
    return 0.0

def check_row_count__79ff096ee36ebe0503f424e88a332ca2(result, expected, **options):
    """Check if file contains the expected row count value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '').strip()
    if not content:
        return 0.0
    expected_count = str(expected.get('expected_count', ''))
    import re
    numbers = re.findall('\\d+', content)
    if expected_count in numbers:
        return 1.0
    for num_str in numbers:
        try:
            num = int(num_str)
            exp = int(expected_count)
            if abs(num - exp) <= 1:
                return 0.5
        except ValueError:
            continue
    return 0.0

def check_row_values__19698d046d0da22c68c3bfcbfbd57b88(result, expected, **options):
    """Check if the row values match expected values with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('label_found'):
        return 0.0
    expected_values = expected.get('expected_values', {})
    tolerance = expected.get('tolerance', 1.0)
    if not expected_values:
        return 0.0
    total_checks = len(expected_values)
    passed = 0
    for (col, exp_val) in expected_values.items():
        actual = result.get('values', {}).get(col)
        if actual is None:
            continue
        try:
            if abs(float(actual) - float(exp_val)) <= tolerance:
                passed += 1
        except (ValueError, TypeError):
            continue
    return passed / total_checks if total_checks > 0 else 0.0

def check_sharper__dd81583a369d4d4452910cdac2a642dc(result, expected, **options):
    """
    Checks that the edited image is sharper than the original while maintaining
    structural similarity.

    result: dict with 'edited_sharpness', 'original_sharpness', 'structural_similarity'
    expected: dict with 'min_similarity' (optional, default 0.5)

    Returns 1.0 if edited image is sharper and structurally similar, 0.0 otherwise.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    edited_sharpness = result.get('edited_sharpness', 0)
    original_sharpness = result.get('original_sharpness', 0)
    structural_similarity = result.get('structural_similarity', 0)
    min_similarity = expected.get('min_similarity', 0.5)
    is_sharper = edited_sharpness > original_sharpness
    is_similar = structural_similarity >= min_similarity
    if is_sharper and is_similar:
        return 1.0
    return 0.0

def check_def_extraction__9f211d2fef2982f0959041bc806f53c4(result, expected, **options):
    """Check that the file contains valid function/class definition lines."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    if not lines:
        return 0.0
    score = 0.0
    if len(lines) >= 1:
        score += 0.3
    def_lines = result.get('def_lines', [])
    non_def_lines = result.get('non_def_lines', [])
    if len(non_def_lines) == 0 and len(def_lines) > 0:
        score += 0.4
    elif len(def_lines) > len(lines) * 0.8:
        score += 0.2
    min_expected = expected.get('min_defs', 5)
    if len(def_lines) >= min_expected:
        score += 0.3
    elif len(def_lines) >= 1:
        score += 0.1
    return min(score, 1.0)

def check_restaurant_names__a68724700bb60550a33fa328934ea4ee(result, expected, **options):
    """Check if restaurant names were entered correctly in A2:A6. Partial credit per name."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_names = result.get('names', [])
    expected_names = expected.get('expected_names', [])
    if not expected_names or not actual_names:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_names)
    for (i, exp_name) in enumerate(expected_names):
        if i < len(actual_names) and actual_names[i] is not None:
            actual_lower = actual_names[i].lower().strip()
            exp_lower = exp_name.lower().strip()
            if exp_lower in actual_lower or actual_lower in exp_lower:
                score += per_item
    return min(score, 1.0)

def check_column_sorted_descending__542fc695553e3e27095b27327032df35(result, expected, **options):
    """Check if date column values are sorted in descending order (newest first)."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', [])
    if not values:
        return 0.0
    str_values = [v for v in values if v is not None]
    if len(str_values) != len(values):
        return 0.0
    expected_sorted = sorted(str_values, reverse=True)
    if str_values == expected_sorted:
        return 1.0
    correct = sum((1 for (a, b) in zip(str_values, expected_sorted) if a == b))
    return correct / len(expected_sorted) * 0.5

def check_iostat_report__d175f88bb1991b0561f10b03c2a108ce(result, expected, **options):
    """Check iostat disk I/O report contains expected headers and keyword patterns.

    Partial credit:
    - 0.5 for containing all expected keywords
    - 0.5 for containing at least min_report_sections distinct iostat iterations (counted by 'avg-cpu' occurrences)
    """
    if not result or result.get('error'):
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    score = 0.0
    keywords = expected.get('keywords', [])
    if keywords:
        all_found = all((kw in content for kw in keywords))
        if all_found:
            score += 0.5
    min_report_sections = int(expected.get('min_report_sections', 1))
    report_count = len(re.findall('avg-cpu', content))
    if report_count >= min_report_sections:
        score += 0.5
    return score

def check_ordered_list__0762adbae4ad235c185f47096cf64c91(result, expected, **options):
    """Check if a list of values matches expected list in order, with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    total = len(expected_values)
    matched = 0
    for (i, exp) in enumerate(expected_values):
        if i < len(actual_values) and actual_values[i].strip() == str(exp).strip():
            matched += 1
    if len(actual_values) > total:
        score = matched / max(total, len(actual_values))
    else:
        score = matched / total
    return min(score, 1.0)

def check_sorted_order__4841f866c3a0d68d763e22faff3652cb(result, expected, **options):
    """Check if papers are sorted by publication date (oldest first).
    Uses partial credit: each correctly placed title earns 0.2.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_titles = result.get('titles', [])
    expected_titles = expected.get('expected_titles', [])
    if not actual_titles or not expected_titles:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_titles)
    for (i, exp_title) in enumerate(expected_titles):
        if i < len(actual_titles):
            actual = str(actual_titles[i]).strip().lower()
            exp = str(exp_title).strip().lower()
            if actual == exp:
                score += per_item
    return min(score, 1.0)

def check_vcard_count__36ae571726da697a61321c69b6ca11f3(result, expected, **options):
    """Check that a vCard file contains the expected number of contacts."""
    if result is None:
        return 0.0
    try:
        count = int(result.strip())
    except (ValueError, AttributeError):
        return 0.0
    min_count = expected.get('min_count', 30)
    return 1.0 if count >= min_count else 0.0

def check_eml_backup__b0f7174902a45f1ac3121d8b78a62ff1(result, expected, **options):
    """Check if expected EML files exist in the backup directory."""
    import re
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    eml_count = result.get('count', 0)
    expected_count = expected.get('expected_count', 2)
    expected_patterns = expected.get('expected_patterns', [])
    score = 0.0
    if eml_count >= expected_count:
        score += 0.5
    if expected_patterns:
        raw_output = result.get('raw_output', '')
        matched = 0
        for pattern in expected_patterns:
            if re.search(pattern, raw_output):
                matched += 1
        if expected_patterns:
            score += 0.5 * (matched / len(expected_patterns))
    elif eml_count >= expected_count:
        score += 0.5
    return min(score, 1.0)

def check_profit_margin_col__60b13be5a7ccd193b76ecf097a198265(result, expected, **options):
    """Check profit margin column header and values."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header', '')
    if header and 'profit' in str(header).lower() and ('margin' in str(header).lower()):
        score += 0.2
    expected_values = expected.get('expected_values', [])
    actual_values = result.get('values', [])
    tolerance = expected.get('tolerance', 0.01)
    if len(actual_values) == len(expected_values):
        correct = 0
        for (actual, exp) in zip(actual_values, expected_values):
            if actual is not None and exp is not None:
                try:
                    if abs(float(actual) - float(exp)) < tolerance:
                        correct += 1
                except (ValueError, TypeError):
                    pass
        score += 0.8 * (correct / len(expected_values))
    return min(score, 1.0)

def check_multi_column_values__9c6205dd9e0fb3145a3021bfaf014d23(result, expected, **options):
    """Check if multiple column cell values match expected values. Supports partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    cells = result.get('cells', {})
    expected_cells = expected.get('expected_cells', {})
    if not expected_cells:
        return 0.0
    total = len(expected_cells)
    matched = 0
    for (cell_ref, expected_val) in expected_cells.items():
        actual_val = cells.get(cell_ref)
        if actual_val is not None and expected_val is not None:
            if actual_val.strip() == str(expected_val).strip():
                matched += 1
    return matched / total if total > 0 else 0.0

def check_sorted_invoices__5d7cb95eaab32c834b0336c7360277e1(result, expected, **options):
    """Check invoices are correctly sorted into matching and discrepant folders.
    Partial credit:
    - 0.33 per correct file in matching folder (2 files = 0.66)
    - 0.34 for correct file in discrepant folder (1 file)
    Penalty for misplaced files.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_matching = result.get('matching_files', [])
    actual_discrepant = result.get('discrepant_files', [])
    expected_matching = expected.get('matching_files', [])
    expected_discrepant = expected.get('discrepant_files', [])
    score = 0.0
    for ef in expected_matching:
        if ef in actual_matching:
            score += 0.33
    for ef in expected_discrepant:
        if ef in actual_discrepant:
            score += 0.34
    wrong_in_matching = [f for f in actual_matching if f not in expected_matching and f.endswith('.pdf')]
    wrong_in_discrepant = [f for f in actual_discrepant if f not in expected_discrepant and f.endswith('.pdf')]
    score -= (len(wrong_in_matching) + len(wrong_in_discrepant)) * 0.2
    return max(0.0, min(score, 1.0))

def check_paper_titles__db9975908e0c2c40a4467f2add6ed446(result, expected, **options):
    """Check if references.bib contains expected paper titles. Partial credit per title."""
    if not result or result.get('error'):
        return 0.0
    content = result.get('content', '').lower()
    if not content:
        return 0.0
    titles = expected.get('titles', [])
    if not titles:
        return 0.0
    score = 0.0
    per_title = 1.0 / len(titles)
    for title in titles:
        title_lower = title.lower()
        key_words = [w for w in title_lower.split() if len(w) > 3]
        if not key_words:
            key_words = title_lower.split()
        matched_words = sum((1 for w in key_words if w in content))
        if len(key_words) > 0 and matched_words / len(key_words) >= 0.7:
            score += per_title
    return min(score, 1.0)

def check_content_align_title_color__63833ad040c5e27e84bd66675fcf45fe(result, expected, **options):
    """Check content alignment and title font color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_align = expected.get('content_alignment', 'CENTER')
    actual_align = result.get('content_alignment', '')
    if actual_align and expected_align.upper() in actual_align.upper():
        score += 0.5
    expected_color = expected.get('title_color', '00FF00')
    actual_color = result.get('title_color', '')
    if actual_color:
        actual_upper = actual_color.upper()
        if expected_color.upper() in actual_upper or actual_upper in ('00FF00', '008000', '00B050', '00B000'):
            score += 0.5
    return min(score, 1.0)

def check_settings_screen_dims__29161befcd93aa27778463661523fc23(result, expected, **options):
    """Check if screen dimensions match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('screen_width') == expected.get('expected_width'):
        score += 0.5
    if result.get('screen_height') == expected.get('expected_height'):
        score += 0.5
    return score

def check_safe_browsing_enabled__949ee9014650e2c21e76ce453947e00b_qw35sft2_7e58fdfd(result, expected, **options):
    """Check if Chrome Safe Browsing is enabled.

    result: "true" or "false" string from enable_safe_browsing getter.
    expected: dict with 'expected' key set to "true".
    Returns 1.0 if safe browsing is enabled, 0.0 otherwise.
    """
    expected_val = expected.get('expected', 'true')
    if isinstance(result, bool):
        result_str = 'true' if result else 'false'
    else:
        result_str = str(result).lower().strip()
    return 1.0 if result_str == expected_val else 0.0

def check_do_not_track__59ed9e9855d0703d98353105912319d0_qw35sft2_eee1b9e4(result, expected, **options):
    """Check if Chrome's Do Not Track setting matches the expected state."""
    if result is None:
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_enabled = expected.get('expected_enabled', True)
    if isinstance(result, bool):
        return 1.0 if result == expected_enabled else 0.0
    if str(result).lower() in ('true', '1', 'yes') and expected_enabled:
        return 1.0
    if str(result).lower() in ('false', '0', 'no') and (not expected_enabled):
        return 1.0
    return 0.0

def check_delta_miles_checkbox__f9b3359b12066a252ed8849a9876d835_qw35sft2_e03606bf(result, expected, **options):
    """Return 1.0 if 'Shop with Miles' checkbox is checked, else 0.0."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_checked = expected.get('checked', True)
    actual_checked = result.get('checked', False)
    return 1.0 if actual_checked == expected_checked else 0.0

def check_delta_trip_type__45facd29054b2deec3804036a06322c9_qw35sft2_35e971f8(result, expected, **options):
    """Check if the trip type is SELECTED as 'One way' (not just present as an option)."""
    if isinstance(result, dict) and 'error' in result and (not result.get('tree_text')):
        return 0.0
    tree_text = result.get('tree_text', '') if isinstance(result, dict) else str(result)
    if not tree_text:
        return 0.0
    expected_type = expected.get('trip_type', 'One way')
    exclude_type = expected.get('exclude_type', 'Round Trip')
    one_way_selected = bool(re.search('one.?way.{0,50}(selected|checked)', tree_text, re.IGNORECASE)) or bool(re.search('(selected|checked).{0,50}one.?way', tree_text, re.IGNORECASE))
    if not one_way_selected:
        return 0.0
    round_trip_selected = bool(re.search('round.?trip.{0,50}(selected|checked)', tree_text, re.IGNORECASE)) or bool(re.search('(selected|checked).{0,50}round.?trip', tree_text, re.IGNORECASE))
    if round_trip_selected:
        return 0.0
    return 1.0

def check_sb_and_dnt__497c76e8a817211059a0ca1d7eafcbe5_qw35sft2_83f45ad6(result, expected, **options):
    """Check Safe Browsing is enabled AND Do Not Track is turned on.

    Partial credit: 0.5 for each sub-goal.
    result: dict with 'safe_browsing' and 'do_not_track' keys ("true"/"false").
    expected: dict with 'safe_browsing' and 'do_not_track' expected values.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('safe_browsing') == expected.get('safe_browsing', 'true'):
        score += 0.5
    if result.get('do_not_track') == expected.get('do_not_track', 'true'):
        score += 0.5
    return score

def check_startup_opens_google__db8335a9df035ba863791bb757e40ea6_qw35sft2_08f31e16(result, expected, **options):
    """
    Check that Chrome is configured to open google.com on startup.
    Requires restore_on_startup == 4 (specific pages) AND startup_urls contains google.com.
    """
    if not isinstance(result, dict):
        return 0.0
    restore = result.get('restore_on_startup', -1)
    startup_urls = result.get('startup_urls') or []
    if restore != 4:
        return 0.0
    has_google = any(('google.com' in str(url).lower() for url in startup_urls))
    has_funbrain = any(('funbrain' in str(url).lower() for url in startup_urls))
    if has_google and (not has_funbrain):
        return 1.0
    return 0.0

def check_appearance_and_dnt__197c5f21c248c0028a67e57d9193addd_qw35sft2_2650bcc9(result, expected, **options):
    """Check that dark mode is off and Do Not Track is enabled."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    color_scheme = result.get('color_scheme', 2) if isinstance(result, dict) else 2
    if color_scheme in [0, 1]:
        score += 0.5
    do_not_track = result.get('do_not_track', False) if isinstance(result, dict) else False
    if do_not_track:
        score += 0.5
    return score

def check_lang_and_dnt__45fdf66a8d514412f8287b18b125d4f9_qw35sft2_4a99185c(result, expected, **options):
    """Check Chrome's interface language and Do Not Track setting with partial credit.
    Score: 0.5 for correct language, 0.5 for correct DNT state.
    """
    if result is None or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    expected_lang = expected.get('expected_lang', '')
    actual_lang = str(result.get('language', '')).lower()
    if expected_lang and (actual_lang.startswith(expected_lang.lower()) or expected_lang.lower() in actual_lang):
        score += 0.5
    expected_dnt = expected.get('expected_dnt', True)
    actual_dnt_str = str(result.get('do_not_track', 'false')).strip().lower()
    actual_dnt = actual_dnt_str == 'true'
    if actual_dnt == expected_dnt:
        score += 0.5
    return score

def check_do_not_track_enabled__ecfeb6992b96a90fa421f7b2969c282f_qw35sft2_e39772e4(result, expected, **options):
    """Check if Chrome's Do Not Track setting is enabled (True)."""
    if result is None:
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if result == 'true' or result is True or result == 1:
        return 1.0
    return 0.0

def check_recreation_title__6368f69cb89cbc42bff053c412aec363_qw35sft2_860c41ff(result, expected, **options):
    """Check if the active tab title contains the expected keyword on a recreation.gov campground page.

    result: dict from active_tab_info with keys: url, title, content
    expected: dict (unwrapped from rules) with keys:
        url_fragment (str): substring expected in URL
        title_keyword (str): keyword expected in title
    """
    if not isinstance(result, dict):
        return 0.0
    url = result.get('url', '') or ''
    title = result.get('title', '') or ''
    url_fragment = expected.get('url_fragment', 'recreation.gov/camping/campgrounds/')
    title_keyword = expected.get('title_keyword', 'Diamond')
    url_ok = url_fragment in url
    title_ok = title_keyword in title
    if url_ok and title_ok:
        return 1.0
    elif url_ok:
        return 0.5
    return 0.0

def check_delta_origin__607589a24fcd7d8000928f466f0dd577_qw35sft2_c52f3ac4(result, expected, **options):
    """Check if origin airport field contains the expected airport code (JFK)."""
    if isinstance(result, dict) and 'error' in result and (not result.get('tree_text')):
        return 0.0
    tree_text = result.get('tree_text', '') if isinstance(result, dict) else str(result)
    if not tree_text:
        return 0.0
    expected_text = expected.get('expected_text', 'JFK')
    alt_text = expected.get('alt_text', 'New York-Kennedy')
    airport_options = [re.escape(expected_text), re.escape(alt_text)] if alt_text else [re.escape(expected_text)]
    airport_pattern = '|'.join(airport_options)
    field_then_airport = re.compile('(?:From|Origin|Departure(?:\\s+city(?:\\s+or\\s+airport)?)?)[^\\n]{0,150}(?:' + airport_pattern + ')', re.IGNORECASE)
    airport_then_field = re.compile('(?:' + airport_pattern + ')[^\\n]{0,100}(?:From|Origin|Departure)', re.IGNORECASE)
    if field_then_airport.search(tree_text):
        return 1.0
    if airport_then_field.search(tree_text):
        return 1.0
    return 0.0

def check_delta_destination_field__696c0a5477f795f89ba4614ef9b530f9_qw35sft2_7ce9bd2d(result, expected, **options):
    """Return 1.0 if the destination field is set to New York City (NYC), else 0.0."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    return 1.0 if result.get('destination_set', False) else 0.0

def check_history_and_dnt__43a7a4d8fb3eedcedaf50f5a89a2ea93_qw35sft2_2f3f6db4(result, expected, **options):
    """Partial credit: 0.5 for YouTube deleted from history, 0.5 for Do Not Track enabled."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    youtube_count = result.get('youtube_count', -1)
    do_not_track = result.get('do_not_track', False)
    if youtube_count == 0:
        score += 0.5
    if do_not_track is True:
        score += 0.5
    return score

def check_enhanced_protection__ba99e87e44fe809095fa3ce0e0316f1d_qw35sft2_938306e7(result, expected, **options):
    """Check if Chrome Enhanced protection Safe Browsing is enabled.

    result: "true" or "false" string from enable_enhanced_safety_browsing getter.
    expected: dict with 'expected' key set to "true".
    Returns 1.0 if enhanced protection is active, 0.0 otherwise.
    """
    expected_val = expected.get('expected', 'true')
    if isinstance(result, bool):
        result_str = 'true' if result else 'false'
    else:
        result_str = str(result).lower().strip()
    return 1.0 if result_str == expected_val else 0.0

def check_contrast_and_saturation__b674a0dc75a50659a17535e66cb88a78_qw35sft2_546b48af(result, expected, **options):
    """Check that both contrast (stddev) and color saturation increased vs original.

    Scoring: 0.5 for contrast increase, 0.5 for saturation increase.
    Original berries.png: avg_std ~51.1, sat_score ~0.239.
    After contrast=45 only: avg_std ~65.7, sat_score ~0.265 (not enough for saturation check).
    Threshold sat_score >= 0.28 requires an explicit saturation boost beyond contrast alone.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    avg_std = result.get('avg_std', 0.0)
    sat_score = result.get('sat_score', 0.0)
    min_std = expected.get('min_contrast_std', 56.0)
    min_sat = expected.get('min_sat_score', 0.28)
    if avg_std >= min_std:
        score += 0.5
    if sat_score >= min_sat:
        score += 0.5
    return score

def check_frame_at_3s__81411ca26f8f6e6edae115d08d9ab086_qw35sft2_01cab48d(result, expected, **options):
    """Check that frame.png pixel-content matches a reference frame extracted at 3s from fullvideo.mp4."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or not result.get('frame_exists', False):
        return 0.0
    user_frame_b64 = result.get('user_frame')
    ref_frame_b64 = result.get('ref_frame')
    if not user_frame_b64 or not ref_frame_b64:
        return 0.0
    try:
        import numpy as np
        from PIL import Image
        user_bytes = base64.b64decode(user_frame_b64)
        ref_bytes = base64.b64decode(ref_frame_b64)
        user_img = Image.open(io.BytesIO(user_bytes)).convert('RGB')
        ref_img = Image.open(io.BytesIO(ref_bytes)).convert('RGB')
        if user_img.size != ref_img.size:
            ref_img = ref_img.resize(user_img.size, Image.LANCZOS)
        user_arr = np.array(user_img, dtype=float)
        ref_arr = np.array(ref_img, dtype=float)
        mae = float(np.mean(np.abs(user_arr - ref_arr)))
        similarity = 1.0 - mae / 255.0
        threshold = expected.get('similarity_threshold', 0.9)
        return 1.0 if similarity >= threshold else 0.0
    except Exception:
        return 0.0

def check_brightness_increased__1dd8caa614c12c2ce7c1b58be9a02c11_qw35sft2_4166e7d3(result, expected, **options):
    """Check if image average brightness meets or exceeds the expected minimum."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or not result.get('file_found', False):
        return 0.0
    avg_brightness = result.get('avg_brightness', 0.0)
    min_brightness = expected.get('min_brightness', 0.0)
    return 1.0 if avg_brightness >= min_brightness else 0.0

def check_grayscale_and_contrast__a80023f11b098859420cd69b0067e271_qw35sft2_8b42bb15(result, expected, **options):
    """Check that image was converted to grayscale mode AND contrast was increased.

    Scoring: 0.5 for grayscale mode ('L'), 0.5 for contrast increase.
    Original berries.png gray_std ~47.4; threshold 50.0 ensures contrast was boosted.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    mode = result.get('mode', '')
    gray_std = result.get('gray_std', 0.0)
    exp_mode = expected.get('expected_mode', 'L')
    min_gray_std = expected.get('min_gray_std', 50.0)
    if mode == exp_mode:
        score += 0.5
    if gray_std >= min_gray_std:
        score += 0.5
    return score

def check_brightness_increased__53b38a939b32dab621a5954331033a3c_qw35sft2_ab8758a4(result, expected, **options):
    """Check if image average brightness meets or exceeds the expected minimum."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or not result.get('file_found', False):
        return 0.0
    avg_brightness = result.get('avg_brightness', 0.0)
    min_brightness = expected.get('min_brightness', 0.0)
    return 1.0 if avg_brightness >= min_brightness else 0.0

def check_scale_and_contrast__4ce9d4657814d8ed93df9fa8b61db4cf_qw35sft2_13cdc41c(result, expected, **options):
    """Check that image was scaled to ~half its original size (640x426) AND contrast increased.

    Original berries.png is 1280x851. Half size = 640x426.
    Scoring: 0.34 for correct width, 0.33 for correct height, 0.33 for contrast increase.
    Original avg_std ~51.1; threshold 56.0 ensures contrast was meaningfully boosted.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    avg_std = result.get('avg_std', 0.0)
    exp_w = expected.get('expected_width', 640)
    exp_h = expected.get('expected_height', 426)
    min_std = expected.get('min_contrast_std', 56.0)
    if abs(width - exp_w) <= 2:
        score += 0.34
    if abs(height - exp_h) <= 2:
        score += 0.33
    if avg_std >= min_std:
        score += 0.33
    return min(score, 1.0)

def check_three_brightness_increased__b3992f4af3f3949be116d5f6a9289f79_qw35sft2_4c2ee848(result, expected, **options):
    """Check brightness increase for squirrel, panda, and heron. Partial credit 1/3 per image."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    checks = [('squirrel', expected.get('min_squirrel', 0.0)), ('panda', expected.get('min_panda', 0.0)), ('heron', expected.get('min_heron', 0.0))]
    for name, threshold in checks:
        img_stat = result.get(name, {})
        if not isinstance(img_stat, dict):
            continue
        if img_stat.get('file_found', False) and img_stat.get('avg_brightness', 0.0) >= threshold:
            score += 1.0 / 3.0
    return min(round(score, 4), 1.0)

def check_freeze_panes__10368a4827622b87edd2c56b2f249e0f_qw35sft2_6266b504(result, expected, **options):
    """Check that freeze_panes is set to freeze row 1 and columns A-B (i.e., 'C2')."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual = result.get('freeze_panes')
    expected_val = expected.get('freeze_panes', 'C2')
    if actual == expected_val:
        return 1.0
    return 0.0

def check_salesrep_jan_total__5238964d8e657c42c8a059231010b871_qw35sft2_f69af80e(result, expected, **options):
    """Check that B12 contains the expected January total value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.5:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_level_secondary__f558bfbb444b8f5736dbe2588385d5fe_qw35sft2_dd3e1b01(result, expected, **options):
    """Check that B8:B18 are all filled with 'Secondary'."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_value = expected.get('level_value', 'Secondary')
    cells = [f'B{r}' for r in range(8, 19)]
    correct = sum((1 for c in cells if result.get(c) == expected_value))
    return correct / len(cells)

def check_income_gross_with_total__9bcff5517fb7493e61233c0a423569f9_qw35sft2_ea0fbb59(result, expected, **options):
    """Check Gross Profit J2:J10 and sum total J11 with partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_gross = result.get('gross_profit', [])
    expected_gross = expected.get('gross_profit', [])
    if actual_gross and expected_gross:
        correct = sum((1 for a, e in zip(actual_gross, expected_gross) if a == e))
        score += 0.5 * (correct / len(expected_gross))
    actual_total = result.get('total')
    expected_total = expected.get('total')
    if actual_total is not None and expected_total is not None:
        if actual_total == expected_total:
            score += 0.5
    return min(score, 1.0)

def check_employee_ages_avg__151640eb3d43a7eb22c551550e103953_qw35sft2_fc1f68c6(result, expected, **options):
    """
    Partial credit:
      0.5 - D2:D29 contains correct ages
      0.5 - E2 contains the average age (within 0.5 tolerance)
    """
    if not result or result.get('error'):
        return 0.0
    dob_values = result.get('dob_values', [])
    age_values = result.get('age_values', [])
    e2_value = result.get('e2_value')
    today = date.today()
    expected_ages = []
    for dob_str in dob_values:
        if dob_str is None:
            expected_ages.append(None)
            continue
        try:
            dob = date.fromisoformat(dob_str)
            expected_age = today.year - dob.year - ((today.month, today.day) < (dob.month, dob.day))
            expected_ages.append(expected_age)
        except Exception:
            expected_ages.append(None)
    correct = 0
    total = 0
    for exp_age, act_age in zip(expected_ages, age_values):
        if exp_age is None:
            continue
        total += 1
        if act_age is not None:
            try:
                if int(round(float(act_age))) == exp_age:
                    correct += 1
            except Exception:
                pass
    age_score = correct / total if total > 0 else 0.0
    valid_ages = [a for a in expected_ages if a is not None]
    expected_avg = sum(valid_ages) / len(valid_ages) if valid_ages else None
    avg_score = 0.0
    if expected_avg is not None and e2_value is not None:
        try:
            if abs(float(e2_value) - expected_avg) <= 0.5:
                avg_score = 1.0
        except Exception:
            pass
    return 0.5 * age_score + 0.5 * avg_score

def check_vlookup_f2_single__dba068db6ac4d383c892aae7e4d7fc62_qw35sft2_c217ff4b(result, expected, **options):
    """Check that cell F2 contains the expected officer name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('F2')
    if actual is None:
        return 0.0
    expected_val = expected.get('F2')
    return 1.0 if str(actual).strip() == str(expected_val).strip() else 0.0

def check_sort_total_row__2cdc28af4ef71e22984f659197334a50_qw35sft2_aa12df02(result, expected, **options):
    """Check sort ascending (D2=442), Total label in A20, and SUM value in D20. Partial credit ~0.33 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    d2 = result.get('d2')
    if d2 == expected.get('expected_d2'):
        score += 0.34
    a20 = result.get('a20')
    if a20 is not None and str(a20).strip().lower() == expected.get('expected_a20', '').lower():
        score += 0.33
    d20 = result.get('d20')
    expected_total = expected.get('expected_d20')
    if d20 is not None and expected_total is not None:
        try:
            if abs(float(d20) - float(expected_total)) < 0.01:
                score += 0.33
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_total_label_and_jan__4e23e64ce56ffa608b1e8fd2866f5ae0_qw35sft2_77da4f1d(result, expected, **options):
    """Check that A12='Total' (0.5) and B12 = Jan sum (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    label = result.get('label')
    if isinstance(label, str) and label.strip() == expected.get('label'):
        score += 0.5
    actual_jan = result.get('jan')
    exp_jan = expected.get('jan')
    if actual_jan is not None and exp_jan is not None:
        try:
            if abs(float(actual_jan) - float(exp_jan)) < 1.0:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return round(score, 4)

def check_first_col_not_empty__b3e437bf5bf6ec56d9f1c4d46601f3ff_qw35sft2_8a9fd9e9(result, expected, **options):
    """Check that column A no longer has an empty header (blank col was deleted)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_header = expected.get('expected_a1_header', 'First Name')
    actual = result.get('a1_header')
    if actual is None:
        return 0.0
    return 1.0 if str(actual).strip() == str(expected_header).strip() else 0.0

def check_monthly_totals__aff1e90d6f0bac6579c8d1279e6c6c3e_qw35sft2_a24f5cbb(result, expected, **options):
    """Check that row 24 has a 'Total' label in B24 and correct SUM values in C24, D24, E24.
    Partial credit: 0.25 per correct cell (label + 3 sums).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    label = result.get('label')
    expected_label = expected.get('expected_label', 'Total')
    if label is not None and str(label).strip().lower() == expected_label.lower():
        score += 0.25
    jan_total = result.get('jan_total')
    expected_jan = expected.get('expected_jan')
    if jan_total is not None and expected_jan is not None:
        try:
            if abs(float(jan_total) - float(expected_jan)) < 0.01:
                score += 0.25
        except (TypeError, ValueError):
            pass
    feb_total = result.get('feb_total')
    expected_feb = expected.get('expected_feb')
    if feb_total is not None and expected_feb is not None:
        try:
            if abs(float(feb_total) - float(expected_feb)) < 0.01:
                score += 0.25
        except (TypeError, ValueError):
            pass
    mar_total = result.get('mar_total')
    expected_mar = expected.get('expected_mar')
    if mar_total is not None and expected_mar is not None:
        try:
            if abs(float(mar_total) - float(expected_mar)) < 0.01:
                score += 0.25
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_weekly_sales_profit_total_row__50f7148385ea6d504119cd2e40f7d1be_qw35sft2_2b548c12(result, expected, **options):
    """Check Profit header, profit values, A12 label, and D12 total."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('d1_header') == expected.get('expected_d1'):
        score += 0.15
    actual_profits = result.get('profit_values', [])
    expected_profits = expected.get('expected_profits', [])
    if actual_profits and expected_profits and (len(actual_profits) == len(expected_profits)):
        correct = sum((1 for a, e in zip(actual_profits, expected_profits) if a is not None and abs(float(a) - float(e)) < 1.0))
        score += 0.35 * (correct / len(expected_profits))
    a12 = result.get('a12_label') or ''
    if str(a12).lower().strip() == expected.get('expected_a12_lower', 'total'):
        score += 0.2
    total = result.get('d12_total')
    expected_total = expected.get('expected_total')
    if total is not None and expected_total is not None:
        if abs(float(total) - float(expected_total)) < 1.0:
            score += 0.3
    return min(score, 1.0)

def check_employee_split_sorted__5d1c2e4438b4aac3fb5db710d29300f5_qw35sft2_ffeb6422(result, expected, **options):
    """Check that data was split correctly AND rows are sorted by Last Name ascending.

    Partial credit:
    - 0.5: All rows have non-null B/C/D values (split was done)
    - 0.5: Rows are sorted by Last Name (column C) in ascending order
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rows = result.get('rows', [])
    if not rows:
        return 0.0
    expected_count = expected.get('expected_row_count', 21)
    filled = [r for r in rows if r.get('first') and r.get('last') and r.get('rank')]
    if len(filled) >= expected_count:
        score += 0.5
    last_names = [r.get('last', '') or '' for r in rows]
    is_sorted = all((last_names[i].lower() <= last_names[i + 1].lower() for i in range(len(last_names) - 1)))
    if is_sorted:
        first_last_expected = expected.get('first_last_name', 'Badman')
        if last_names and last_names[0] == first_last_expected:
            score += 0.5
    return min(score, 1.0)

def check_c1_format_d1_fix__3be2ab82396c7145b9caded82bed3999_qw35sft2_d313d34b(result, expected, **options):
    """Check D1 uses TEXT for 2 decimal places AND C1 is formatted to 3 decimal places.
    Partial credit: 0.5 for D1 TEXT formula with correct value, 0.5 for C1 format.
    """
    if result.get('error'):
        return 0.0
    score = 0.0
    formula = result.get('d1_formula', '') or ''
    value = result.get('d1_value', '') or ''
    c1_fmt = result.get('c1_number_format', '') or ''
    if 'TEXT' in formula.upper():
        score += 0.25
    expected_value = expected.get('expected_d1_value', 'The price is 19.50 euros.')
    if value.strip() == expected_value:
        score += 0.25
    expected_fmt = expected.get('expected_c1_format', '0.000')
    if c1_fmt == expected_fmt:
        score += 0.5
    return score

def check_period_rate_max_in_d1__14298c6976685b9e8d10b60e8490521a_qw35sft2_e583f219(result, expected, **options):
    """Check Period Rate header, max-value green font, and max period rate value stored in D1."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('C1') == expected.get('header'):
        score += 0.33
    if result.get('C20_font_color', '').lower() == expected.get('max_color', '00ff00').lower():
        score += 0.33
    exp_max = expected.get('max_value', 14.724)
    d1_val = result.get('D1_value')
    if d1_val is not None:
        try:
            if abs(float(d1_val) - exp_max) < 0.05:
                score += 0.34
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_sort_and_sum__8b987ab7e49940ce5b75cd3329aaf643_qw35sft2_0cb12aa2(result, expected, **options):
    """Check that data is sorted ascending by Date Time and G1 contains total quantity (780)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_sorted_ascending') is True:
        score += 0.5
    g1 = result.get('g1_value')
    expected_sum = expected.get('expected_total', 780)
    if g1 is not None:
        try:
            if abs(float(g1) - float(expected_sum)) < 0.5:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_maturity_sorted__6a0b6ecdb454d8389c50fe9e251b1580_qw35sft2_806b62f8(result, expected, **options):
    """
    Check that:
    1. C1 = 'Maturity Date' (0.4 pts)
    2. Data rows C2:C10 are sorted ascending by maturity date (0.4 pts)
    3. C2 is the earliest maturity date (0.2 pts)
    expected keys: header_expected, earliest_date
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header_c', '')
    expected_header = expected.get('header_expected', 'Maturity Date')
    if isinstance(header, str) and header.strip().lower() == expected_header.strip().lower():
        score += 0.4
    if result.get('is_sorted') and result.get('c2_value') is not None:
        score += 0.4
    expected_earliest = expected.get('earliest_date', '')
    c2 = result.get('c2_value', '')
    if c2 and expected_earliest and (c2 == expected_earliest):
        score += 0.2
    return min(score, 1.0)

def check_passfail_and_validation__1231c66fffc58f9761299238f8444bbb_qw35sft2_b3126219(result, expected, **options):
    """Check D2:D29 are filled correctly (Pass>=50, Fail<50, Held=--) and data validation exists."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_values = expected.get('expected_d_values', {})
    d_values = result.get('d_values', {})
    if expected_values:
        correct = sum((1 for k, v in expected_values.items() if d_values.get(k) == v))
        score += 0.5 * (correct / len(expected_values))
    if result.get('validation_correct'):
        score += 0.5
    elif result.get('has_list_validation'):
        score += 0.25
    return min(score, 1.0)

def check_salesrep_label__d6d9e5e1d2ca3d967d1969fd1a1e0c2c_qw35sft2_7b6ca218(result, expected, **options):
    """Check that A12 contains the expected label (case-insensitive strip)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    label = result.get('label')
    expected_label = expected.get('label', 'Total')
    if label is None:
        return 0.0
    if isinstance(label, str) and label.strip().lower() == expected_label.strip().lower():
        return 1.0
    return 0.0

def check_seqno_and_total_row__f9ba38ec0e00ec28c9c0ef5057c79916_qw35sft2_0d24a241(result, expected, **options):
    """Check Seq No. (B2:B29), 'Total:' label in D30, and total sales value in E30."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    seq_nos = result.get('seq_nos', [])
    expected_seq = [f'No. {i}' for i in range(1, 29)]
    if seq_nos == expected_seq:
        score += 0.34
    d30 = result.get('d30')
    expected_label = expected.get('expected_d30', 'Total:')
    if d30 is not None and str(d30).strip() == expected_label:
        score += 0.33
    e30 = result.get('e30')
    expected_sum = expected.get('expected_e30')
    if e30 is not None and expected_sum is not None:
        try:
            if abs(float(e30) - float(expected_sum)) < 0.01:
                score += 0.33
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_employee_ages__add589b113a1be1eec737c1ab22661fd_qw35sft2_2484c826(result, expected, **options):
    """Check that D2:D29 contains correct ages computed from DOB in column C."""
    if not result or result.get('error'):
        return 0.0
    dob_values = result.get('dob_values', [])
    age_values = result.get('age_values', [])
    if not dob_values or not age_values:
        return 0.0
    today = date.today()
    correct = 0
    total = 0
    for dob_str, age_val in zip(dob_values, age_values):
        if dob_str is None:
            continue
        total += 1
        try:
            dob = date.fromisoformat(dob_str)
            expected_age = today.year - dob.year - ((today.month, today.day) < (dob.month, dob.day))
            if age_val is not None:
                actual_age = int(round(float(age_val)))
                if actual_age == expected_age:
                    correct += 1
        except Exception:
            pass
    if total == 0:
        return 0.0
    return correct / total

def check_vlookup_f2_f4_rows__22ff441442432fb67bd17ef3eb97a292_qw35sft2_eb420e13(result, expected, **options):
    """Check that F2, F3, F4 each contain the correct officer name (partial credit)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    for cell in ('F2', 'F3', 'F4'):
        actual = result.get(cell)
        expected_val = expected.get(cell)
        if actual is not None and str(actual).strip() == str(expected_val).strip():
            score += 1.0 / 3
    return min(score, 1.0)

def check_income_net_sales_gross__c6e3aa9fd1cbfd34b4473b9e8f31e349_qw35sft2_80a50e1c(result, expected, **options):
    """Check Net Sales (E2:E10) and Gross Profit (J2:J10) with partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_net = result.get('net_sales', [])
    expected_net = expected.get('net_sales', [])
    if actual_net and expected_net:
        correct = sum((1 for a, e in zip(actual_net, expected_net) if a == e))
        score += 0.5 * (correct / len(expected_net))
    actual_gross = result.get('gross_profit', [])
    expected_gross = expected.get('gross_profit', [])
    if actual_gross and expected_gross:
        correct = sum((1 for a, e in zip(actual_gross, expected_gross) if a == e))
        score += 0.5 * (correct / len(expected_gross))
    return min(score, 1.0)

def check_level_primary_range__74bebda311b0ec1306e448288c71ed13_qw35sft2_10d2d467(result, expected, **options):
    """Check that B3:B6 are all filled with 'Primary'."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_value = expected.get('level_value', 'Primary')
    cells = ['B3', 'B4', 'B5', 'B6']
    correct = sum((1 for c in cells if result.get(c) == expected_value))
    return correct / len(cells)

def check_total_row_state__5ac69261c8fb294b16f18049ece06ce9_qw35sft2_8f27b526(result, expected, **options):
    """Check that the Total row label and all 6 monthly sums are correct."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    weight_per_item = 1.0 / 7
    label = result.get('label')
    if isinstance(label, str) and label.strip() == expected.get('label'):
        score += weight_per_item
    month_keys = ['jan', 'feb', 'mar', 'apr', 'may', 'jun']
    for key in month_keys:
        actual = result.get(key)
        exp_val = expected.get(key)
        if actual is not None and exp_val is not None:
            try:
                if abs(float(actual) - float(exp_val)) < 1.0:
                    score += weight_per_item
            except (TypeError, ValueError):
                pass
    return round(min(score, 1.0), 4)

def check_single_row_hidden__93bd58adf830a704a8880938e35d28c4_qw35sft2_d3296ba3(result, expected, **options):
    """Return 1.0 if the target row is hidden, 0.0 otherwise."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    return 1.0 if result.get('hidden', False) else 0.0

def check_weekly_sales_sorted_by_profit__99a00a677b41fb8db78ec79d5381f0b2_qw35sft2_fc61000e(result, expected, **options):
    """Check Profit column header, profit values present, and data sorted descending by profit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('d1_header') == expected.get('expected_header'):
        score += 0.2
    actual_profits = result.get('profit_values', [])
    expected_set = set((round(float(v), 0) for v in expected.get('expected_profits_set', [])))
    if actual_profits:
        actual_set = set((round(float(p), 0) for p in actual_profits if p is not None))
        if actual_set == expected_set:
            score += 0.3
    profits_clean = [float(p) for p in actual_profits if p is not None]
    if len(profits_clean) >= 2:
        is_sorted = all((profits_clean[i] >= profits_clean[i + 1] for i in range(len(profits_clean) - 1)))
        if is_sorted:
            score += 0.5
    return min(score, 1.0)

def check_feb_max__663a803da15ba1a41afdeb0114613f5e_qw35sft2_31440d0b(result, expected, **options):
    """Check that cell D24 contains the correct maximum for Feb sales (980)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_value = expected.get('expected_value')
    actual_value = result.get('value')
    if actual_value is None:
        return 0.0
    try:
        actual_num = float(actual_value)
        expected_num = float(expected_value)
        return 1.0 if abs(actual_num - expected_num) < 0.01 else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_d1_usd_fix__62d484d2d5eca5e0b713e4dbdcc095e8_qw35sft2_dcdf4e3f(result, expected, **options):
    """Check D1 shows price with 2 decimal places and currency changed to USD.
    Partial credit: 0.5 for TEXT formula, 0.5 for 'USD' in displayed value with 2 decimals.
    """
    if result.get('error'):
        return 0.0
    score = 0.0
    formula = result.get('d1_formula', '') or ''
    value = result.get('d1_value', '') or ''
    if 'TEXT' in formula.upper():
        score += 0.5
    must_contain_decimal = expected.get('must_contain_decimal', '19.50')
    must_contain_currency = expected.get('must_contain_currency', 'USD')
    if must_contain_decimal in value and must_contain_currency in value:
        score += 0.5
    return score

def check_ramp_accel_diff__e64652ba3bf9a74b525e231bf3f391d2_qw35sft2_dbfad9cd(result, expected, **options):
    """Check the Acceleration Difference column (E = B - D) and column B/D fill.

    Partial credit (0.20 each):
    - E1 header contains 'diff' (case-insensitive)
    - E2 is approximately 2.957 (B2 - D2 = -2.24047 - (-5.19723))
    - E30 is approximately 2.92 (B30 - D30)
    - B30 is approximately 1.70 (column B filled correctly)
    - D30 is approximately -1.22 (column D filled correctly)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tol = 0.02
    score = 0.0
    e1 = result.get('e1') or ''
    if 'diff' in e1.lower():
        score += 0.2
    e2 = result.get('e2')
    try:
        if e2 is not None and abs(float(e2) - 2.957) <= tol:
            score += 0.2
    except (TypeError, ValueError):
        pass
    e30 = result.get('e30')
    try:
        if e30 is not None and abs(float(e30) - 2.92) <= tol:
            score += 0.2
    except (TypeError, ValueError):
        pass
    b30 = result.get('b30')
    try:
        if b30 is not None and abs(float(b30) - 1.7) <= 0.015:
            score += 0.2
    except (TypeError, ValueError):
        pass
    d30 = result.get('d30')
    try:
        if d30 is not None and abs(float(d30) - -1.22) <= tol:
            score += 0.2
    except (TypeError, ValueError):
        pass
    return round(score, 4)

def check_period_rate_sum_c26__fd409a83b18f09cf45613c5173377d4d_qw35sft2_cb0acd0d(result, expected, **options):
    """Check Period Rate header, max-value row green font, and sum in C26."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('C1') == expected.get('header'):
        score += 0.25
    exp_max = expected.get('max_value', 14.724)
    c20_val = result.get('C20_value')
    if c20_val is not None and abs(float(c20_val) - exp_max) < 0.05:
        score += 0.25
    if result.get('C20_font_color', '').lower() == expected.get('max_color', '00ff00').lower():
        score += 0.25
    exp_sum = expected.get('sum_value', 87.883)
    c26_val = result.get('C26_value')
    if c26_val is not None:
        try:
            if abs(float(c26_val) - exp_sum) / exp_sum < 0.01:
                score += 0.25
        except (TypeError, ValueError, ZeroDivisionError):
            pass
    return min(score, 1.0)

def check_sort_and_sum__7b14f7b3fd242abdf7c663f7ccb97562_qw35sft2_d1f6f816(result, expected, **options):
    """Check sort ascending (D2=442) and SUM formula in D20 (=76079). Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    d2 = result.get('d2')
    if d2 == expected.get('expected_d2'):
        score += 0.5
    d20 = result.get('d20')
    expected_total = expected.get('expected_d20')
    if d20 is not None and expected_total is not None:
        try:
            if abs(float(d20) - float(expected_total)) < 0.01:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score

def check_sort_and_sumif__bda50c4eee48b4e3a15799921e84da62_qw35sft2_5d609264(result, expected, **options):
    """Check that data is sorted ascending and G1 contains total amazon.com quantity (388)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_sorted_ascending') is True:
        score += 0.5
    g1 = result.get('g1_value')
    expected_sumif = expected.get('expected_sumif', 388)
    if g1 is not None:
        try:
            if abs(float(g1) - float(expected_sumif)) < 0.5:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_sales_total__c61f5322a84515967a3c4cc16d8ae654_qw35sft2_4f7f37c0(result, expected, **options):
    """Check that the Sales column total equals the expected sum."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_total = expected.get('expected_total', 56478)
    actual = result.get('sum_value')
    if actual is None:
        return 0.0
    try:
        return 1.0 if abs(float(actual) - float(expected_total)) < 0.01 else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_maturity_total__461252d0c8044977d30ecc240f2dc9cd_qw35sft2_677d3c70(result, expected, **options):
    """
    Check that:
    1. C1 = 'Maturity Date' header added (0.34 pts)
    2. A11 = expected total label (0.33 pts)
    3. B11 = 990 (total loan days) (0.33 pts)
    expected keys: header_expected, a11_expected, b11_expected
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header_c', '')
    expected_header = expected.get('header_expected', 'Maturity Date')
    if isinstance(header, str) and header.strip().lower() == expected_header.strip().lower():
        score += 0.34
    a11 = result.get('a11', '')
    expected_a11 = expected.get('a11_expected', 'Total')
    if isinstance(a11, str) and a11.strip().lower() == expected_a11.strip().lower():
        score += 0.33
    b11 = result.get('b11')
    expected_b11 = expected.get('b11_expected', 990)
    try:
        if b11 is not None and int(b11) == int(expected_b11):
            score += 0.33
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_passfail_and_count__9a44e82214177e638fd33343975b4139_qw35sft2_22da825c(result, expected, **options):
    """Check D2:D29 correct values, E1 COUNTIF Pass result, and data validation."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_d_values = expected.get('expected_d_values', {})
    d_values = result.get('d_values', {})
    if expected_d_values:
        correct = sum((1 for k, v in expected_d_values.items() if d_values.get(k) == v))
        score += 0.34 * (correct / len(expected_d_values))
    expected_pass_count = expected.get('pass_count', 11)
    e1_val = result.get('e1_value')
    e2_val = result.get('e2_value')
    e1_matches = False
    e2_matches = False
    if e1_val is not None:
        try:
            e1_matches = int(float(e1_val)) == expected_pass_count
        except (ValueError, TypeError):
            pass
    if e2_val is not None:
        try:
            e2_matches = int(float(e2_val)) == expected_pass_count
        except (ValueError, TypeError):
            pass
    if e1_matches or e2_matches:
        score += 0.33
    if result.get('has_list_validation'):
        score += 0.33
    return min(score, 1.0)

def check_row_freeze__ae87df329e6b7ba3f255d9a2acdfeda6_qw35sft2_0ffacc85(result, expected, **options):
    """Check that only row 1 is frozen (freeze_panes == 'A2'), not columns."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual = result.get('freeze_panes')
    expected_val = expected.get('freeze_panes', 'A2')
    if actual == expected_val:
        return 1.0
    return 0.0

def check_seqno_and_sales_sum__00713006369c01f94f764c7eb6008543_qw35sft2_1ee1692c(result, expected, **options):
    """Check Seq No. column (B2:B29) is 'No. 1'..'No. 28' and E30 equals expected total sales."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    seq_nos = result.get('seq_nos', [])
    expected_seq = [f'No. {i}' for i in range(1, 29)]
    if seq_nos == expected_seq:
        score += 0.5
    e30 = result.get('e30')
    expected_sum = expected.get('expected_sum')
    if e30 is not None and expected_sum is not None:
        try:
            if abs(float(e30) - float(expected_sum)) < 0.01:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score

def check_salesrep_last3_totals__ce77d41e3a7b3efa48e4655c22c8aa27_qw35sft2_5bfb12c9(result, expected, **options):
    """Check Apr/May/Jun totals in E12:G12. Equal weight (1/3 each)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    checks = [('apr', 'apr_total'), ('may', 'may_total'), ('jun', 'jun_total')]
    for result_key, expected_key in checks:
        actual = result.get(result_key)
        exp_val = expected.get(expected_key)
        if actual is not None and exp_val is not None:
            try:
                if abs(float(actual) - float(exp_val)) < 0.5:
                    score += 1.0 / 3.0
            except (TypeError, ValueError):
                pass
    return min(score, 1.0)

def check_student_c3_c4__87a420fb18a4347285a9615e2d7a9d87_qw35sft2_1102acc6(result, expected, **options):
    """Check that C3 and C4 both contain the expected student name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_name = expected.get('student_name', 'Blake Dreary')
    correct = sum((1 for c in ['C3', 'C4'] if result.get(c) == expected_name))
    return correct / 2

def check_vlookup_f2_f12_all__202132c1158925d79d1ba222174d8f66_qw35sft2_934b09f4(result, expected, **options):
    """Check all 11 officer name cells F2:F12 contain correct values. Partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    cells = ('F2', 'F3', 'F4', 'F5', 'F6', 'F7', 'F8', 'F9', 'F10', 'F11', 'F12')
    score = 0.0
    for cell in cells:
        actual = result.get(cell)
        expected_val = expected.get(cell)
        if actual is not None and str(actual).strip() == str(expected_val).strip():
            score += 1.0 / len(cells)
    return min(score, 1.0)

def check_income_net_sales_and_cost__43eee44b990a73c4b5a6e2bc138833a7_qw35sft2_b852ddbd(result, expected, **options):
    """Check Net Sales (E2:E10) and Total Cost (I2:I10) with partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_net = result.get('net_sales', [])
    expected_net = expected.get('net_sales', [])
    if actual_net and expected_net:
        correct = sum((1 for a, e in zip(actual_net, expected_net) if a == e))
        score += 0.5 * (correct / len(expected_net))
    actual_cost = result.get('total_cost', [])
    expected_cost = expected.get('total_cost', [])
    if actual_cost and expected_cost:
        correct = sum((1 for a, e in zip(actual_cost, expected_cost) if a == e))
        score += 0.5 * (correct / len(expected_cost))
    return min(score, 1.0)

def check_save_and_transition__5e93f9ddd16cd129b875ad7386e8a5e5_qw35sft2_17efed89(result, expected, **options):
    """Check that pre.pptx was saved on Desktop (0.5) and all slides have fade transition (0.5)."""
    if isinstance(result, str) or result is None:
        return 0.0
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
    transitions = result.get('all_transitions', [])
    slide_count = result.get('slide_count', 0)
    expected_transition = expected.get('transition_type', 'fade')
    if slide_count > 0 and transitions:
        all_have_transition = all((isinstance(t, str) and t.lower() == expected_transition.lower() for t in transitions))
        if all_have_transition:
            score += 0.5
    return min(score, 1.0)

def check_presenter_console_disabled__c6ff7e494d134a42d6e5c420e00bf4fd_qw35sft2_502d7a7b(result, expected, **options):
    """Check if presenter console is disabled. Returns 1.0 if disabled, 0.0 otherwise."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    return 1.0 if result.get('presenter_disabled', False) else 0.0

def check_stretch_and_no_picture3__bd30ea01dba8f440569b065c94e4f32f_qw35sft2_1c1836e7(result, expected, **options):
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tolerance_cm = 0.8
    slide_width = result.get('slide_width_cm', 25.4)
    slide_height = result.get('slide_height_cm', 19.05)
    image_width = result.get('image_width_cm', 0.0)
    image_height = result.get('image_height_cm', 0.0)
    image_left = result.get('image_left_cm', -99.0)
    image_top = result.get('image_top_cm', -99.0)
    fills_width = abs(image_width - slide_width) <= tolerance_cm
    fills_height = abs(image_height - slide_height) <= tolerance_cm
    image_fills_page = fills_width or fills_height
    if fills_width:
        expected_top = (slide_height - image_height) / 2
        centered = abs(image_left) <= tolerance_cm and abs(image_top - expected_top) <= tolerance_cm
    elif fills_height:
        expected_left = (slide_width - image_width) / 2
        centered = abs(image_top) <= tolerance_cm and abs(image_left - expected_left) <= tolerance_cm
    else:
        centered = False
    if image_fills_page and centered:
        score += 0.5
    if not result.get('picture3_exists', True):
        score += 0.5
    return min(score, 1.0)

def check_strikethrough_and_transition__b374fc122f888120a74e40bee1da5133_qw35sft2_09c1ad4a(result, expected, **options):
    """Check strikethrough on two Finance Meetings items + fade transition on slide 5. Partial credit."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    strike_state = result.get('strike_state', {})
    transition_type = result.get('transition_type', 'none')
    strikethrough_targets = expected.get('strikethrough_targets', [])
    if strikethrough_targets:
        per_item = 0.67 / len(strikethrough_targets)
        for text in strikethrough_targets:
            if strike_state.get(text) is True:
                score += per_item
    expected_transition = expected.get('transition_type', '')
    if expected_transition and expected_transition.lower() in str(transition_type).lower():
        score += 0.33
    return min(round(score, 4), 1.0)

def check_stretch_and_fade_transition__676dd2418e7e3a84c9083aa5ad9f76e2_qw35sft2_15654b86(result, expected, **options):
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tolerance_cm = 0.8
    slide_width = result.get('slide_width_cm', 25.4)
    slide_height = result.get('slide_height_cm', 19.05)
    image_width = result.get('image_width_cm', 0.0)
    image_height = result.get('image_height_cm', 0.0)
    image_left = result.get('image_left_cm', -99.0)
    image_top = result.get('image_top_cm', -99.0)
    fills_width = abs(image_width - slide_width) <= tolerance_cm
    fills_height = abs(image_height - slide_height) <= tolerance_cm
    image_fills_page = fills_width or fills_height
    if fills_width:
        expected_top = (slide_height - image_height) / 2
        centered = abs(image_left) <= tolerance_cm and abs(image_top - expected_top) <= tolerance_cm
    elif fills_height:
        expected_left = (slide_width - image_width) / 2
        centered = abs(image_top) <= tolerance_cm and abs(image_left - expected_left) <= tolerance_cm
    else:
        centered = False
    if image_fills_page and centered:
        score += 0.5
    transition = result.get('transition_type', 'none').lower()
    if 'fade' in transition:
        score += 0.5
    return min(score, 1.0)

def check_picture_heights__eff65ebb8ed6d102db81f31da1a823ed_qw35sft2_de0e2541(result, expected, **options):
    """Check picture heights on slides 3, 4, 6. Partial credit 1/3 per slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tolerance = 0.5
    score = 0.0
    slide3_actual = result.get('slide3_height_cm')
    slide3_target = expected.get('slide3_height_cm')
    if slide3_actual is not None and abs(slide3_actual - slide3_target) <= tolerance:
        score += 1.0 / 3.0
    slide4_actual = result.get('slide4_height_cm')
    slide4_target = expected.get('slide4_height_cm')
    if slide4_actual is not None and abs(slide4_actual - slide4_target) <= tolerance:
        score += 1.0 / 3.0
    slide6_actual = result.get('slide6_height_cm')
    slide6_target = expected.get('slide6_height_cm')
    if slide6_actual is not None and abs(slide6_actual - slide6_target) <= tolerance:
        score += 1.0 / 3.0
    return min(round(score, 4), 1.0)

def check_strikethrough_multi__7f3881fcd58bcba327d08f20ed190dcc_qw35sft2_e7c6ddca(result, expected, **options):
    """Check strikethrough on multiple bullets across sections with partial credit."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    targets = expected.get('targets', [])
    if not targets:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(targets)
    for text in targets:
        if result.get(text) is True:
            score += per_item
    return min(round(score, 4), 1.0)

def check_titlecase_and_center__c181651f58b43b69af608583a0523891_qw35sft2_81f55bca(result, expected, **options):
    """Check title case (0.5) and title paragraph center alignment (0.5).
    expected is already the rules dict (unwrapped by get_rule()).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_case_applied') == expected.get('title_case_applied', True):
        score += 0.5
    if result.get('title_centered') == expected.get('title_centered', True):
        score += 0.5
    return score

def check_word_colors__336f1c4245e680e1aace133fb243059a_qw35sft2_e8073d66(result, expected, **options):
    """Check if specified words have expected font colors. Returns partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    word_colors = result.get('word_colors', {}) if isinstance(result, dict) else {}
    checks = expected.get('word_color_checks', [])
    tolerance = expected.get('color_tolerance', 50)
    if not checks:
        return 0.0
    color_match = lambda actual, target, tol: len((actual or '').upper().lstrip('#')) == 6 and len((target or '').upper().lstrip('#')) == 6 and all((abs(int((actual or '').upper().lstrip('#')[i * 2:(i + 1) * 2], 16) - int((target or '').upper().lstrip('#')[i * 2:(i + 1) * 2], 16)) <= tol for i in range(3)))
    correct = sum((1 for c in checks if color_match(word_colors.get(c.get('word', '').lower()), c.get('color', ''), tolerance)))
    return correct / len(checks)

def check_first_two_double__c985596bc3953d61c7d274fc30ba4960_qw35sft2_b823e19e(result, expected, **options):
    """Check that the first two content paragraphs both have double line spacing.
    Partial credit: 0.5 per paragraph.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('para0_spacing') == expected.get('para0_spacing', 'double'):
        score += 0.5
    if result.get('para2_spacing') == expected.get('para2_spacing', 'double'):
        score += 0.5
    return score

def check_word_colors__d89e455c3ca4dbc17e773411cf8f66df_qw35sft2_d3435f5f(result, expected, **options):
    """Check if specified words have expected font colors. Returns partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    word_colors = result.get('word_colors', {}) if isinstance(result, dict) else {}
    checks = expected.get('word_color_checks', [])
    tolerance = expected.get('color_tolerance', 50)
    if not checks:
        return 0.0
    correct = 0
    for c in checks:
        actual = (word_colors.get(c.get('word', '').lower()) or '').upper().lstrip('#')
        target = (c.get('color', '') or '').upper().lstrip('#')
        if len(actual) == 6 and len(target) == 6:
            try:
                r1, g1, b1 = (int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16))
                r2, g2, b2 = (int(target[0:2], 16), int(target[2:4], 16), int(target[4:6], 16))
                if abs(r1 - r2) <= tolerance and abs(g1 - g2) <= tolerance and (abs(b1 - b2) <= tolerance):
                    correct += 1
            except ValueError:
                pass
    return correct / len(checks)

def check_title_alignment__6435e9b69d0bbcb863c89964e3084688_qw35sft2_ed3834f4(result, expected, **options):
    """Check whether the document title paragraph has the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('alignment', 'CENTER').upper()
    actual_alignment = (result.get('alignment') or 'LEFT').upper()
    return 1.0 if actual_alignment == expected_alignment else 0.0

def check_word_colors__e97799f0add4268fcc31cdd1e95ac277_qw35sft2_e161288d(result, expected, **options):
    """Check if specified words have expected font colors. Returns partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    word_colors = result.get('word_colors', {}) if isinstance(result, dict) else {}
    checks = expected.get('word_color_checks', [])
    tolerance = expected.get('color_tolerance', 50)
    if not checks:
        return 0.0
    correct = 0
    for c in checks:
        actual = (word_colors.get(c.get('word', '').lower()) or '').upper().lstrip('#')
        target = (c.get('color', '') or '').upper().lstrip('#')
        if len(actual) == 6 and len(target) == 6:
            try:
                r1, g1, b1 = (int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16))
                r2, g2, b2 = (int(target[0:2], 16), int(target[2:4], 16), int(target[4:6], 16))
                if abs(r1 - r2) <= tolerance and abs(g1 - g2) <= tolerance and (abs(b1 - b2) <= tolerance):
                    correct += 1
            except ValueError:
                pass
    return correct / len(checks)

def check_mixed_spacing__ffe17a8bd50575f09eec01a68fed60c6_qw35sft2_e0536b7e(result, expected, **options):
    """Check that first paragraph has double spacing and second paragraph has 1.5 spacing.
    Partial credit: 0.5 per paragraph matching expected spacing.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('para0_spacing') == expected.get('para0_spacing', 'double'):
        score += 0.5
    if result.get('para2_spacing') == expected.get('para2_spacing', '1.5'):
        score += 0.5
    return score

def check_titlecase_and_doublespace__6e2b5ade87cb088089f67dc5070eb77b_qw35sft2_867c1751(result, expected, **options):
    """Check title case (0.5) and double line spacing on all content paragraphs (0.5).
    expected is already the rules dict (unwrapped by get_rule()).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_case_applied') == expected.get('title_case_applied', True):
        score += 0.5
    if result.get('all_double_spacing') == expected.get('all_double_spacing', True):
        score += 0.5
    return score

def check_word_colors__3d9f327a66b9025c03e34cafe8c88fe3_qw35sft2_fcfc98f6(result, expected, **options):
    """Check if specified words have expected font colors. Returns partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    word_colors = result.get('word_colors', {}) if isinstance(result, dict) else {}
    checks = expected.get('word_color_checks', [])
    tolerance = expected.get('color_tolerance', 50)
    if not checks:
        return 0.0
    correct = 0
    for c in checks:
        actual = (word_colors.get(c.get('word', '').lower()) or '').upper().lstrip('#')
        target = (c.get('color') or '').upper().lstrip('#')
        if len(actual) == 6 and len(target) == 6:
            try:
                r1, g1, b1 = (int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16))
                r2, g2, b2 = (int(target[0:2], 16), int(target[2:4], 16), int(target[4:6], 16))
                if abs(r1 - r2) <= tolerance and abs(g1 - g2) <= tolerance and (abs(b1 - b2) <= tolerance):
                    correct += 1
            except ValueError:
                pass
    return correct / len(checks)

def check_settings_window_size__a8bc55c18275bf61543bfdfa6bb3d447_qw35sft2_068a7c7a(result, expected, **options):
    """Check WIDTH and HEIGHT in settings.py match expected values (partial credit)."""
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('width') == expected.get('width'):
        score += 0.5
    if result.get('height') == expected.get('height'):
        score += 0.5
    return score

def check_dir_exists__e66cc6b0ec6479dc303b558a9cc72642_qw35sft2_43067696(result, expected, **options):
    """Check if a directory exists on the VM."""
    if not isinstance(result, dict):
        return 0.0
    should_exist = expected.get('should_exist', True)
    exists = result.get('exists', False)
    return 1.0 if exists == should_exist else 0.0

def check_res_txt_descending__da27b9fa71ba67953d0b9f649a3197fa_qw35sft2_7fd0b028(result, expected, **options):
    """Check that res.txt contains the sorted array in descending order."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    expected_numbers = expected.get('expected_numbers', '90 64 34 25 22 12 11')
    content_normalized = ' '.join(content.split())
    expected_normalized = ' '.join(expected_numbers.split())
    if expected_normalized and expected_normalized in content_normalized:
        return 1.0
    return 0.0

def check_invoiceGES_in_problematic__cea86e2d544f9d987b2adcf034f36c4a_qw35sft2_6981b395(result, expected, **options):
    """Return 1.0 if Invoice # GES-20220215-82.pdf is inside Desktop/problematic folder."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('folder_exists'):
        score += 0.4
    if result.get('file_in_problematic'):
        score += 0.6
    return min(score, 1.0)

def check_2017_2018_cities__1272197382489180193ce8b35860bb11_qw35sft2_f2b45d1d(result, expected, **options):
    """Check 2017 and 2018 conference city cells with equal partial credit (1/6 each)."""
    if not result or result.get('error'):
        return 0.0
    ALIASES = {'long beach': ['long beach', 'los angeles', 'la'], 'montreal': ['montreal', 'montréal']}
    keys = ['c14', 'c15', 'c16', 'c17', 'c18', 'c19']
    score = 0.0
    per_item = 1.0 / len(keys)
    for key in keys:
        actual = result.get(key)
        exp = expected.get(key, '')
        matched = False
        if actual and exp:
            norm_actual = str(actual).strip().lower()
            norm_exp = str(exp).strip().lower()
            if norm_actual == norm_exp:
                matched = True
            else:
                for aliases in ALIASES.values():
                    if norm_exp in aliases and norm_actual in aliases:
                        matched = True
                        break
        if matched:
            score += per_item
    return min(score, 1.0)

def check_small_compression__98edc79a225cc1dd6b5d8a23724241ab_qw35sft2_cc37c421(result, expected, **options):
    """Check that the JPEG file exists and is under the strict max_size threshold."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') or result.get('size_bytes', -1) < 0:
        return 0.0
    size_bytes = result.get('size_bytes', -1)
    max_size = expected.get('max_size', 204800)
    if size_bytes > 0 and size_bytes < max_size:
        return 1.0
    return 0.0

def check_answers_tests23__ed93db9128130b7a782216ac4507a0ca_qw35sft2_80825339(result, expected, **options):
    """Partial credit: 0.5 for test 2 answer, 0.5 for test 3 answer in Answer.docx."""
    if result.get('error'):
        return 0.0
    full_text = result.get('full_text', '')
    score = 0.0
    if expected.get('test2_answer', 'baaad') in full_text:
        score += 0.5
    if expected.get('test3_answer', 'aaaaa') in full_text:
        score += 0.5
    return score

def check_sar_disk_report__fc8ac0599dbb02ca56427b3ae265c796_qw35sft2_a50d1235(result, expected, **options):
    """Check that Disk_IO_Report.txt exists and contains disk I/O data."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    min_lines = expected.get('min_lines', 5)
    if result.get('line_count', 0) >= min_lines:
        score += 0.5
    if result.get('has_disk_data', False):
        score += 0.5
    return min(score, 1.0)

def check_txt_multihop_gemini__61501dfdaa17fec1d3c116cb7f744d5c_qw35sft2_3940d7cb(result, expected, **options):
    """Check that gemini_multihop.txt has content from Multi_Hop GEMINI responses."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    required_substrings = expected.get('required_substrings', [])
    min_char_count = expected.get('min_char_count', 100)
    if result.get('char_count', 0) >= min_char_count:
        score += 0.5
    content = result.get('content', '')
    if required_substrings:
        matches = sum((1 for s in required_substrings if s in content))
        score += 0.5 * (matches / len(required_substrings))
    elif result.get('contains_iliad', False):
        score += 0.5
    return min(score, 1.0)

def check_imdb_top250_nav__7c977607bff83aa73fb51b31d1f513df_qw35sft2_0c77bee4(result, expected, **options):
    """Check if the active tab URL contains the expected IMDB Top 250 path."""
    if result is None:
        return 0.0
    url = result if isinstance(result, str) else str(result)
    if 'error' in url.lower():
        return 0.0
    expected_fragment = expected.get('expected_url', 'imdb.com/chart/top')
    return 1.0 if expected_fragment in url else 0.0

def check_dir_exists__e103d8977cd2fefeec7972ff8169e36d_qw35sft2_b561c678(result, expected, **options):
    """Check if the directory existence matches expected state."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_exists = expected.get('expected_exists', True)
    return 1.0 if result.get('exists') == expected_exists else 0.0

def check_settings_snake_size__85e39531da2c848a1afaf10e25ba3dad_qw35sft2_a4edccc9(result, expected, **options):
    """Check that SNAKE_SIZE in settings.py matches expected value."""
    if result.get('error'):
        return 0.0
    if result.get('snake_size') == expected.get('snake_size'):
        return 1.0
    return 0.0

def check_invoiceTII_in_problematic__7ecb63a4609903641df39f4754b7248f_qw35sft2_b90e9695(result, expected, **options):
    """Return 1.0 if invoice TII-20220301-90.pdf is inside Desktop/problematic folder."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('folder_exists'):
        score += 0.4
    if result.get('file_in_problematic'):
        score += 0.6
    return min(score, 1.0)

def check_book_copied_to_documents__55d3e622f74c4c6d9051e6e358a64ecd_qw35sft2_b3db2c69(result, expected, **options):
    """Check that the book PDF was copied to the Documents folder."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('file_exists', False) else 0.0

def check_bubblesort_impl_and_output__8c65eee451a1bb106a548123af80ea08_qw35sft2_3002e435(result, expected, **options):
    """Partial credit: 0.5 for bubbleSort.py having real implementation, 0.5 for res.txt correct output."""
    score = 0.0
    py_content = result.get('py_content', '') or ''
    txt_content = result.get('txt_content', '') or ''
    has_loop = 'for ' in py_content or 'while ' in py_content
    no_todo_only = '#TODO' not in py_content or (has_loop and '#TODO' in py_content)
    has_comparison = '>' in py_content or '<' in py_content
    has_real_body = has_loop and has_comparison and ('def bubbleSort' in py_content)
    if has_real_body:
        score += 0.5
    expected_output = expected.get('expected_output', '11 12 22 25 34 64 90')
    content_normalized = ' '.join(txt_content.split())
    expected_normalized = ' '.join(expected_output.split())
    if expected_normalized and expected_normalized in content_normalized:
        score += 0.5
    return min(score, 1.0)

def check_ext_and_resize__3001e2e091ccfaef8fccec23c5f15501_qw35sft2_ad209965(result, expected, **options):
    """Partial-credit metric: Lisp extension installed (0.5) + image resized (0.5).

    Args:
        result:   dict from get_combined_state__09914c54 with keys
                  'extensions', 'img_width', 'img_height'
        expected: rules dict with 'width' and 'height' (default 128)

    Returns:
        float in [0.0, 1.0]
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_output = result.get('extensions', '')
    for line in ext_output.split('\n'):
        if 'lisp' in line.strip().lower():
            score += 0.5
            break
    expected_w = expected.get('width', 128)
    expected_h = expected.get('height', 128)
    if result.get('img_width') == expected_w and result.get('img_height') == expected_h:
        score += 0.5
    return min(score, 1.0)

def check_sampled_conf_cities__fb22e54de099b00f4d6a127b12703f11_qw35sft2_b8b70a1e(result, expected, **options):
    """Check 3 sampled conference city cells with partial credit (0.33+0.34+0.33)."""
    if not result or result.get('error'):
        return 0.0
    checks = [('c3', expected.get('c3', ''), 0.33), ('c13', expected.get('c13', ''), 0.34), ('c20', expected.get('c20', ''), 0.33)]
    score = 0.0
    for key, exp_val, weight in checks:
        actual = result.get(key)
        if actual and exp_val and (str(actual).strip().lower() == str(exp_val).strip().lower()):
            score += weight
    return min(score, 1.0)

def check_tally_book_amount_sum__289ae9cd25fb72b631d3d5c97c4b9f82_qw35sft2_fdfcb7c8(result, expected, **options):
    """Check PDF receipt saved with correct name AND tally book total matches expected value.
    Partial credit: 0.5 per sub-goal.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('pdf_exists', False):
        score += 0.5
    actual_total = result.get('total') if isinstance(result, dict) else None
    if actual_total is not None:
        expected_total = expected.get('expected_total', 22.18)
        tolerance = expected.get('tolerance', 0.05)
        if abs(actual_total - expected_total) <= tolerance:
            score += 0.5
    return score

def check_paper03_and_year__d508d903f596b0ecb03360bdda6892d4_qw35sft2_f08f2902(result, expected, **options):
    """Partial credit: 0.5 for paper03.pdf existing, 0.5 for correct pub year in paper03_year.txt."""
    if result is None or isinstance(result, str):
        return 0.0
    score = 0.0
    if result.get('pdf_exists'):
        score += 0.5
    expected_year = str(expected.get('expected_year', '2017'))
    actual_year = result.get('year_content', '').strip()
    if actual_year == expected_year or expected_year in actual_year:
        score += 0.5
    return score

def check_goodbye_world__f65942d684aa73fa8a727494a6335877_qw35sft2_c83e2311(result, expected, **options):
    """Check that main.py contains a goodbye_world function that prints 'Goodbye, world!'.

    Partial credit:
      0.5 - function definition 'def goodbye_world' is present
      0.5 - correct print statement with 'Goodbye, world!' is present
    """
    if not result or not isinstance(result, str):
        return 0.0
    content = result
    score = 0.0
    if 'def goodbye_world' in content:
        score += 0.5
    if 'Goodbye, world!' in content:
        score += 0.5
    return score

def check_sar_cpu_report__5dfc721833abb9e182bb18218c19d630_qw35sft2_1e2c5a0d(result, expected, **options):
    """Check that System_Resources_Report.txt exists with at least 30 CPU data lines."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    min_cpu_lines = expected.get('min_cpu_lines', 30)
    if result.get('cpu_lines', 0) >= min_cpu_lines:
        score += 0.7
    if result.get('has_header', False):
        score += 0.3
    return min(score, 1.0)

def check_answers_and_count__25c4c78ad235f6f7dec2a3ce1f7f2c6f_qw35sft2_30c1ab5f(result, expected, **options):
    """Partial credit: 0.25 each for test 2 answer, test 3 answer in docx;
    0.5 for correct count value in xlsx cell B9."""
    full_text = result.get('full_text', '')
    cell_value = result.get('cell_value')
    score = 0.0
    if expected.get('test2_answer', 'baaad') in full_text:
        score += 0.25
    if expected.get('test3_answer', 'aaaaa') in full_text:
        score += 0.25
    expected_cell = expected.get('cell_value', 3)
    try:
        if cell_value is not None and int(float(str(cell_value))) == int(expected_cell):
            score += 0.5
    except (ValueError, TypeError):
        pass
    return round(min(score, 1.0), 4)

def check_settings_fps__751d0a9ef75d8ec2c4c27e57248ad37c_qw35sft2_daa09262(result, expected, **options):
    """Check that FPS in settings.py matches expected value."""
    if result.get('error'):
        return 0.0
    if result.get('fps') == expected.get('fps'):
        return 1.0
    return 0.0

def check_invoice243729_in_problematic__da5edbd81a2d596fb7eb6ef7b52c2151_qw35sft2_2c597180(result, expected, **options):
    """Return 1.0 if Invoice # 243729.pdf is inside Desktop/problematic folder."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('folder_exists'):
        score += 0.4
    if result.get('file_in_problematic'):
        score += 0.6
    return min(score, 1.0)

def check_clipboard_path__241466d2f7566981ed09035f56716bb4_qw35sft2_6fab5826(result, expected, **options):
    """Check that clipboard contains exactly the expected file path."""
    if result is None or result.get('error'):
        return 0.0
    clipboard = result.get('clipboard', '').strip()
    expected_path = expected.get('expected_path', '').strip()
    if not expected_path:
        return 0.0
    if clipboard == expected_path:
        return 1.0
    if expected_path in clipboard:
        return 1.0
    return 0.0

def check_conda_install__d5abdc9092a51546a82222c972d5edf6_qw35sft2_d9f2a9fa(result, expected, **options):
    """
    Check conda is installed AND .bashrc has been configured for auto-activation.
    0.5 for conda binary present, 0.5 for .bashrc conda init block present.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('conda_installed', False):
        score += 0.5
    if result.get('bashrc_configured', False):
        score += 0.5
    return min(score, 1.0)

def check_icml_cities__e73f742a9bd3e72f68e6861886fdd0c7_qw35sft2_25ed0820(result, expected, **options):
    """Check all 7 ICML city cells with equal partial credit (1/7 each)."""
    if not result or result.get('error'):
        return 0.0
    keys = ['c3', 'c6', 'c9', 'c12', 'c15', 'c18', 'c21']
    score = 0.0
    per_item = 1.0 / len(keys)
    for key in keys:
        actual = result.get(key)
        exp = expected.get(key, '')
        if not actual or not exp:
            continue
        norm_actual = str(actual).strip().lower()
        norm_exp = str(exp).strip().lower()
        matched = norm_actual == norm_exp
        if not matched:
            for aliases in _ICML_CITY_ALIASES_qw35sft2_45c2e8.values():
                if norm_exp in aliases and norm_actual in aliases:
                    matched = True
                    break
        if matched:
            score += per_item
    return min(score, 1.0)

def check_paper01_and_count__101069bcaae5b37861b743495ae4859d_qw35sft2_03bcee0a(result, expected, **options):
    """Partial credit: 0.5 for paper01.pdf existing, 0.5 for correct paper count in paper_count.txt."""
    if result is None or isinstance(result, str):
        return 0.0
    score = 0.0
    if result.get('pdf_exists'):
        score += 0.5
    expected_count = str(expected.get('expected_count', '5'))
    actual_count = result.get('count_content', '').strip()
    if actual_count == expected_count:
        score += 0.5
    return score

def check_res_txt_sorted__68c8947721f1221211aba42cd6d1735d_qw35sft2_c967d941(result, expected, **options):
    """Check that res.txt contains the expected sorted array output."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    expected_output = expected.get('expected_output', '')
    content_normalized = ' '.join(content.split())
    expected_normalized = ' '.join(expected_output.split())
    if expected_normalized and expected_normalized in content_normalized:
        return 1.0
    expected_nums = expected.get('expected_numbers', [])
    if expected_nums:
        nums_str = ' '.join((str(n) for n in expected_nums))
        content_nums = ' '.join(content.split())
        if nums_str in content_nums:
            return 1.0
    return 0.0

def check_tally_book_last_row__97f6f4530ffb08062fee8d0568f59984_qw35sft2_d1e6ad53(result, expected, **options):
    """Check PDF extracted to receipts folder and tally book last row has correct service, month, amount."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('pdf_found'):
        score += 0.5
    if not result.get('error'):
        expected_service = expected.get('expected_service', 'AWS')
        expected_month = expected.get('expected_month', 2023.12)
        expected_amount = expected.get('expected_amount', 10.02)
        if result.get('service') == expected_service:
            score += 0.17
        actual_month = result.get('month')
        if actual_month is not None and abs(actual_month - expected_month) < 0.001:
            score += 0.165
        actual_amount = result.get('amount')
        if actual_amount is not None and abs(actual_amount - expected_amount) < 0.05:
            score += 0.165
    return min(score, 1.0)

def check_copy_move_state__ec373221258728f3163857f2bdfd353a_qw35sft2_a2d7b604(result, expected, **options):
    """0.25 each: dir1/file1, dir2/file1, dir3/file1, original file1 gone."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir1_has_file1'):
        score += 0.25
    if result.get('dir2_has_file1'):
        score += 0.25
    if result.get('dir3_has_file1'):
        score += 0.25
    if result.get('original_removed'):
        score += 0.25
    return min(round(score, 6), 1.0)

def check_clock_dual_settings__443802668a0bfacbf0757782ecb2ad43_qw35sft2_e19209e8(result, expected, **options):
    """Check both clock-format and clock-show-weekday with partial credit (0.5 each).

    expected (already unwrapped by get_rule()):
        target_format   - expected substring in clock-format output, e.g. "'12h'"
        target_weekday  - expected substring in clock-show-weekday output, e.g. "true"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    target_format = expected.get('target_format', "'12h'")
    target_weekday = expected.get('target_weekday', 'true')
    clock_format = result.get('clock_format', '') if isinstance(result, dict) else ''
    clock_weekday = result.get('clock_show_weekday', '') if isinstance(result, dict) else ''
    if target_format in clock_format:
        score += 0.5
    if target_weekday in clock_weekday:
        score += 0.5
    return score

def check_screen_blank_timeout__40046a4f8d5906bbddcb5abf1fee6b07_qw35sft2_1cac0dec(result, expected, **options):
    """Check if screen blank timeout matches expected idle-delay value."""
    if not isinstance(result, dict):
        return 0.0
    expected_delay = expected.get('expected_idle_delay', 600)
    return 1.0 if result.get('idle_delay') == expected_delay else 0.0

def check_screen_lock_enabled__74e051f75a4f9998d114f6407d55ddd5_qw35sft2_c77e69c9(result, expected, **options):
    """Check if automatic screen lock is enabled (lock-enabled = true)."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    lock_str = result.get('lock_enabled', '')
    return 1.0 if 'true' in lock_str.lower() else 0.0

def check_volume_level__d6b3f95410b1cde5a46558c01aaafc46_qw35sft2_4d0db765(result, expected, **options):
    """Check that sink volume is within tolerance of target_percent."""
    import re
    if not isinstance(result, str) or not result:
        return 0.0
    target = expected.get('target_percent', 100)
    tolerance = expected.get('tolerance', 5)
    matches = re.findall('(\\d+)%', result)
    if not matches:
        return 0.0
    actual = int(matches[0])
    return 1.0 if abs(actual - target) <= tolerance else 0.0

def check_utc0_24h_clock__c55229b53381e6587f722d00658cbd02_qw35sft2_68eb61e3(result, expected, **options):
    """Check that system timezone is UTC+0 AND GNOME clock is set to 24-hour format.

    Partial credit scoring:
      - 0.5 for timezone offset == '+0000'
      - 0.5 for clock_format containing '24h'

    Args:
        result: dict with keys 'timezone_offset' and 'clock_format'
        expected: dict (from rules) – checked keys are embedded in logic
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    score = 0.0
    offset = str(result.get('timezone_offset', '')).strip()
    if offset == '+0000':
        score += 0.5
    clock_fmt = str(result.get('clock_format', '')).strip().strip('\'"')
    if clock_fmt == '24h':
        score += 0.5
    return score

def check_copy_4dirs__4f3015d5aa890d8fea505e363d3f7aff_qw35sft2_13a61b9c(result, expected, **options):
    """0.25 each: file1 in dir1, dir2, dir3, and dir4."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir1_has_file1'):
        score += 0.25
    if result.get('dir2_has_file1'):
        score += 0.25
    if result.get('dir3_has_file1'):
        score += 0.25
    if result.get('dir4_has_file1'):
        score += 0.25
    return min(round(score, 6), 1.0)

def check_rename_and_sibling__9cc29f7146497c596e333ae0ebaf5848_qw35sft2_964736b0(result, expected, **options):
    """
    Partial-credit check:
      0.5 - Desktop folder renamed from todo_list_Jan_1 to todo_list_Jan_2
      0.5 - New folder todo_list_Jan_3 created on Desktop
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('renamed') and result.get('old_absent', True):
        score += 0.5
    if result.get('created'):
        score += 0.5
    return score

def check_notif_clock__accf3abf7a1d21f7f68eefd98414a494_qw35sft2_cd16c09b(result, expected, **options):
    """Check DND enabled (0.5) and clock set to 24h format (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if 'false' in result.get('show_banners', ''):
        score += 0.5
    if '24h' in result.get('clock_format', ''):
        score += 0.5
    return score

def check_power_dim_and_blank__4a865a2b73390cd1be0bc5f929d03ff5_qw35sft2_fc5c977b(result, expected, **options):
    """Check idle-dim is false (0.5) and idle-delay is 0/Never (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('idle_dim', '').lower() == 'false':
        score += 0.5
    if result.get('idle_delay') == 0:
        score += 0.5
    return score

def check_screen_lock_with_delay__8b9cedf68783a02fb2242a963784d1c2_qw35sft2_5da34c49(result, expected, **options):
    """Check screen lock enabled (0.5) and blank screen delay matches 600s / 10 min (0.5)."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    lock_str = result.get('lock_enabled', '')
    if 'true' in lock_str.lower():
        score += 0.5
    delay_str = result.get('idle_delay', '')
    expected_delay = str(expected.get('idle_delay', ''))
    if expected_delay and expected_delay in delay_str:
        score += 0.5
    return score

def check_php_stats__a63cdece1fd90a60a3e737af233f1764_qw35sft2_aba72e41(result, expected, **options):
    """Partial-credit check for PHP file count and total line count.

    Scoring:
        0.5 - ~/php_file_count.txt contains the expected file count
        0.5 - ~/php_line_count.txt contains the expected line count
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_file_count = str(expected.get('expected_file_count', '4'))
    expected_line_count = str(expected.get('expected_line_count', '54'))
    file_count_content = str(result.get('file_count', ''))
    line_count_content = str(result.get('line_count', ''))
    if expected_file_count in file_count_content:
        score += 0.5
    if expected_line_count in line_count_content:
        score += 0.5
    return score

def check_move_failed_notebooks__96695aa3604bded58ed2ec7f3c72c8de_qw35sft2_199fbbbb(result, expected, **options):
    """Partial credit: 0.5 for files in ./fails, 0.5 for originals deleted."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    fails_files = result.get('fails_files', '')
    orig_files = result.get('orig_failed_remaining', '')
    required_in_fails = expected.get('required_in_fails', [])
    if required_in_fails and all((f in fails_files for f in required_in_fails)):
        score += 0.5
    if not orig_files:
        score += 0.5
    return score

def check_output_and_backup__9fff5b0d3da3f15287b691b4b1944f6b_qw35sft2_939b4b84(result, expected, **options):
    """Check output.txt and /tmp/output_backup.txt both have correct <br/>-appended content.

    Scoring:
      0.5 - output.txt contains '1<br/>', '2<br/>', '3<br/>'
      0.5 - /tmp/output_backup.txt contains the same lines
    """
    if not isinstance(result, dict):
        return 0.0
    expected_lines = expected.get('expected_lines', ['1<br/>', '2<br/>', '3<br/>'])
    score = 0.0
    output_txt = result.get('output_txt') or ''
    if all((s in output_txt for s in expected_lines)):
        score += 0.5
    backup_txt = result.get('backup_txt') or ''
    if all((s in backup_txt for s in expected_lines)):
        score += 0.5
    return score

def check_rename_and_move__eb200f79b699eb5d233b5379c7e1ea0e_qw35sft2_8153b443(result, expected, **options):
    """
    Partial-credit check:
      0.5 - todo_list_Jan_2 exists in home directory
      0.5 - original todo_list_Jan_1 no longer on Desktop
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('in_home'):
        score += 0.5
    if result.get('old_gone'):
        score += 0.5
    return score

def check_volume_level__2ac92884cbfeefbd5a2d6929c20839d7_qw35sft2_5e6ad00c(result, expected, **options):
    """Check that sink volume is within tolerance of target_percent."""
    import re
    if not isinstance(result, str) or not result:
        return 0.0
    target = expected.get('target_percent', 100)
    tolerance = expected.get('tolerance', 5)
    matches = re.findall('(\\d+)%', result)
    if not matches:
        return 0.0
    actual = int(matches[0])
    return 1.0 if abs(actual - target) <= tolerance else 0.0

def check_screen_blank_never__3696e962dd7f0719e2ddde3289beabec_qw35sft2_89043873(result, expected, **options):
    """Check if screen blank is set to the expected idle-delay value."""
    if not isinstance(result, dict):
        return 0.0
    expected_delay = expected.get('expected_idle_delay', 0)
    return 1.0 if result.get('idle_delay') == expected_delay else 0.0

def check_restore_and_copy__f767d224a0cdea04b11256d9cffabd41_qw35sft2_49933203(result, expected, **options):
    """Partial credit: 0.5 for poster restored to Desktop, 0.5 for copy saved in Documents."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('file_on_desktop', False):
        score += 0.5
    if result.get('file_in_documents', False):
        score += 0.5
    return min(score, 1.0)

def check_copy_and_count_fails__8f62d3b2706e8f98a78e22aea189146a_qw35sft2_24b70f63(result, expected, **options):
    """Partial credit: 0.5 for files copied to ./fails, 0.5 for count.txt with correct count."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    fails_files = result.get('fails_files', '')
    count_txt = result.get('count_txt', '')
    required_files = expected.get('required_files', [])
    if required_files and all((f in fails_files for f in required_files)):
        score += 0.5
    expected_count = str(expected.get('expected_count', ''))
    if expected_count and expected_count in count_txt:
        score += 0.5
    return score

def check_notif_dnd_lockscreen__7a1296a554f735dccd7fa86e30d68dbd_qw35sft2_9d57f044(result, expected, **options):
    """Check DND enabled (0.5) and lock-screen notifications disabled (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if 'false' in result.get('show_banners', ''):
        score += 0.5
    if 'false' in result.get('show_in_lock_screen', ''):
        score += 0.5
    return score

def check_utc0_ntp_disabled__397a38a7528d157681d1e57316ff4073_qw35sft2_008be2a1(result, expected, **options):
    """Check that system timezone is UTC+0 AND automatic time sync (NTP) is disabled.

    Partial credit scoring:
      - 0.5 for timezone offset == '+0000'
      - 0.5 for NTP status == 'no' (disabled)

    Args:
        result: dict with keys 'timezone_offset' and 'ntp_status'
        expected: dict (from rules) with keys 'expected_offset' and 'expected_ntp'
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    score = 0.0
    offset = str(result.get('timezone_offset', '')).strip()
    if offset == '+0000':
        score += 0.5
    ntp_status = str(result.get('ntp_status', '')).strip().lower()
    if ntp_status == 'no':
        score += 0.5
    return score

def check_screen_lock_with_delay__b5802fb719515ed2e2b308492efceb55_qw35sft2_3bd8d49f(result, expected, **options):
    """Check screen lock enabled (0.5) and blank screen delay matches 120s / 2 min (0.5)."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    lock_str = result.get('lock_enabled', '')
    if 'true' in lock_str.lower():
        score += 0.5
    delay_str = result.get('idle_delay', '')
    expected_delay = str(expected.get('idle_delay', ''))
    if expected_delay and expected_delay in delay_str:
        score += 0.5
    return score

def check_subdir_split_permissions__31b8299b7936d108e594f17a62518506_qw35sft2_16fc09a7(result, expected, **options):
    """Check subDir1 files are 600 (0.5 pts) and subDir2 files are 644 (0.5 pts)."""
    if not result:
        return 0.0
    score = 0.0
    subdir1_perms = result.get('subdir1_permissions', [])
    subdir2_perms = result.get('subdir2_permissions', [])
    expected_subdir1 = expected.get('expected_subdir1_permission', '600')
    expected_subdir2 = expected.get('expected_subdir2_permission', '644')
    if subdir1_perms and all((p == expected_subdir1 for p in subdir1_perms)):
        score += 0.5
    if subdir2_perms and all((p == expected_subdir2 for p in subdir2_perms)):
        score += 0.5
    return score

def check_volume_settings__1201dde5af3c06b515ed16eae4fad685_qw35sft2_279ea4a6(result, expected, **options):
    """Partial credit: 0.5 for max volume (100%), 0.5 for Settings app open."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if '100%' in result.get('volume_output', ''):
        score += 0.5
    if result.get('settings_open', False):
        score += 0.5
    return score

def check_restore_and_rename__ee992b5c4f2de4a7b0ef0f22ec1a67b1_qw35sft2_a576718e(result, expected, **options):
    """Partial credit: 0.5 for restoring poster to Desktop, 0.5 for renaming to party_poster.webp."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    has_new = result.get('has_new_name', False)
    has_old = result.get('has_old_name', False)
    score = 0.0
    if has_new or has_old:
        score += 0.5
    if has_new and (not has_old):
        score += 0.5
    return min(score, 1.0)

def check_notif_screenlock__0ee0cafb23f5579534516156a3b19db3_qw35sft2_0abda510(result, expected, **options):
    """Check DND enabled (0.5) and screen auto-lock disabled (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if 'false' in result.get('show_banners', ''):
        score += 0.5
    if 'false' in result.get('lock_enabled', ''):
        score += 0.5
    return score

def check_power_triple_state__570bc607f27286d53f7b1a8eaf2af624_qw35sft2_ef529bf8(result, expected, **options):
    """Check three power settings: dim off (0.33), blank never (0.34), auto-suspend on (0.33)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('idle_dim', '').lower() == 'false':
        score += 0.33
    if result.get('idle_delay') == 0:
        score += 0.34
    if result.get('sleep_inactive_ac_type', '').lower() == 'suspend':
        score += 0.33
    return min(score, 1.0)

def check_screen_lock_with_delay__5e9c7d265f840de83da743ff562a1457_qw35sft2_7bb5e860(result, expected, **options):
    """Check screen lock enabled (0.5) and blank screen delay matches 60s / 1 min (0.5)."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    lock_str = result.get('lock_enabled', '')
    if 'true' in lock_str.lower():
        score += 0.5
    delay_str = result.get('idle_delay', '')
    expected_delay = str(expected.get('idle_delay', ''))
    if expected_delay and expected_delay in delay_str:
        score += 0.5
    return score

def check_output_and_count__bb3776631f95d3c4e2bcd96582b451ff_qw35sft2_a7da54b2(result, expected, **options):
    """Check output.txt has <br/> appended lines and count.txt has correct line count.

    Scoring:
      0.5 - output.txt contains '1<br/>', '2<br/>', '3<br/>'
      0.5 - count.txt contains the string '3'
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    output_txt = result.get('output_txt') or ''
    if all((s in output_txt for s in ['1<br/>', '2<br/>', '3<br/>'])):
        score += 0.5
    count_txt = result.get('count_txt') or ''
    if count_txt.strip() == '3':
        score += 0.5
    return score

def check_restore_and_wallpaper__354747c46b7fabc52ac1f7b0a9e657e8_qw35sft2_0291f617(result, expected, **options):
    """Partial credit: 0.5 for restoring poster to Desktop, 0.5 for wallpaper set to that file."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('file_on_desktop', False):
        score += 0.5
    wallpaper_uri = result.get('wallpaper_uri', '')
    if 'poster_party_night' in wallpaper_uri:
        score += 0.5
    return min(score, 1.0)

def check_notif_sounds__139fa3103c161a5284bcac4b5dc4f632_qw35sft2_2a8c221a(result, expected, **options):
    """Check DND enabled (0.5) and event sounds disabled (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if 'false' in result.get('show_banners', ''):
        score += 0.5
    if 'false' in result.get('event_sounds', ''):
        score += 0.5
    return score

def check_power_dim_and_suspend__6cd1c7feeea3fc52512c1552e8d5b3b2_qw35sft2_c03caca2(result, expected, **options):
    """Check idle-dim is false (0.5) and auto-suspend is enabled/suspend (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('idle_dim', '').lower() == 'false':
        score += 0.5
    if result.get('sleep_inactive_ac_type', '').lower() == 'suspend':
        score += 0.5
    return score

def check_copy_and_rename__44eff79b79182e6eeb0c7debb1bc9f8c_qw35sft2_0bd73cb5(result, expected, **options):
    """0.25 each: dir1/file1, dir2/file1, dir3/file1, file1.bak in home."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir1_has_file1'):
        score += 0.25
    if result.get('dir2_has_file1'):
        score += 0.25
    if result.get('dir3_has_file1'):
        score += 0.25
    if result.get('home_has_file1_bak'):
        score += 0.25
    return min(round(score, 6), 1.0)

def check_archive_content__7b7da38700f679a895053728c5a7d35b_qw35sft2_c00fd3bb(result, expected, **options):
    """Partial credit: 0.34 archive exists, 0.33 old_file1.txt in archive, 0.33 old_file2.txt in archive."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    file1 = expected.get('file1', 'old_file1.txt')
    file2 = expected.get('file2', 'old_file2.txt')
    contents = result.get('contents', [])
    if result.get('archive_exists'):
        score += 0.34
    if any((file1 in c for c in contents)):
        score += 0.33
    if any((file2 in c for c in contents)):
        score += 0.33
    return round(min(score, 1.0), 2)

def check_two_dir_notebook_split__e71563318b46b13cadfcb76e6b9ce246_qw35sft2_e2a4d5be(result, expected, **options):
    """Partial credit: 0.5 for *failed.ipynb in ./fails, 0.5 for non-failed in ./passing."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    fails_files = result.get('fails_files', '')
    passing_files = result.get('passing_files', '')
    required_fails = expected.get('required_fails', [])
    if required_fails and all((f in fails_files for f in required_fails)):
        score += 0.5
    required_passing = expected.get('required_passing', [])
    if required_passing and all((f in passing_files for f in required_passing)):
        score += 0.5
    return score

def check_accessibility_two_toggles__8ee8727a5fa2032ce96fa89a79ac9fab_qw35sft2_da39f4c7(result, expected, **options):
    """Check large-text and high-contrast toggles with partial credit (0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('large_text') is True:
        score += 0.5
    if result.get('high_contrast') is True:
        score += 0.5
    return score

def check_eml_backup__d81f9d153146f886b82c4131d09d1ccb_qw35sft2_7f674302(result, expected, **options):
    """
    Partial-credit check for EML backup:
      - 0.5 if the number of .eml files equals the expected count
      - 0.5 if the expected subject fragment is found in the directory listing
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    try:
        actual_count = int(str(result.get('count', '0')).strip())
        target_count = int(expected.get('count', 2))
        if actual_count == target_count:
            score += 0.5
    except (ValueError, AttributeError, TypeError):
        logger_qw35sft2_d0992a.warning('Could not parse EML count from result: %s', result.get('count'))
    listing = result.get('listing', '')
    subject_fragment = expected.get('subject_fragment', '')
    if subject_fragment and subject_fragment in listing:
        score += 0.5
    return score

def check_bills_starred_and_important__ed9097e7c75920408b4536024b8da2a4_qw35sft2_53fd849d(result, expected, **options):
    """Partial credit: 0.5 for all Bills emails starred, 0.5 for all tagged as Important."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('message_count', 0) == 0:
        return 0.0
    score = 0.0
    if result.get('all_starred'):
        score += 0.5
    if result.get('all_important'):
        score += 0.5
    return score

def check_bills_triple_combo__d21b0cc63afcca51c02fdd16b5d862e0_qw35sft2_a12fd2c7(result, expected, **options):
    """Partial credit: 0.34 starred + 0.33 unread + 0.33 Important tag."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('message_count', 0) == 0:
        return 0.0
    score = 0.0
    if result.get('all_starred'):
        score += 0.34
    if result.get('all_unread'):
        score += 0.33
    if result.get('all_important'):
        score += 0.33
    return min(score, 1.0)

def check_eml_count__2731b9abd5cfbad9ed4df8aae737addc_qw35sft2_70f9d6b7(result, expected, **options):
    """Check that the number of .eml files matches the expected count."""
    try:
        actual = int(str(result).strip())
        target = int(expected.get('count', 0))
        return 1.0 if actual == target else 0.0
    except (ValueError, AttributeError, TypeError):
        logger_qw35sft2_dce5f0.error('Could not parse EML count from result: %s', result)
        return 0.0

def check_bills_starred_unread__8e8eeb1588f1109e98ccb02a0faa787c_qw35sft2_e7b5f644(result, expected, **options):
    """Partial credit: 0.5 for all-starred, 0.5 for all-unread."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('message_count', 0) == 0:
        return 0.0
    score = 0.0
    if result.get('all_starred'):
        score += 0.5
    if result.get('all_unread'):
        score += 0.5
    return score

def check_eml_subject__dfeb48225188ee18fb4de9d6f0048829_qw35sft2_631026b2(result, expected, **options):
    """Check that a specific subject fragment appears in the EML backup directory listing."""
    if not result or not isinstance(result, str):
        return 0.0
    subject_fragment = expected.get('subject_fragment', '')
    if not subject_fragment:
        return 1.0 if '.eml' in result else 0.0
    return 1.0 if subject_fragment in result else 0.0

def check_bills_starred_and_receipts__0b7b1b3c91e02a8222608f02952fc214_qw35sft2_aa546db1(result, expected, **options):
    """Partial credit: 0.5 for all Bills emails starred, 0.5 for Receipts folder created."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('all_starred'):
        score += 0.5
    if result.get('receipts_folder_exists'):
        score += 0.5
    return score

def check_bills_starred_and_filter__44c5673d772aeaf1fd4e6dcc8a32111d_qw35sft2_624b5785(result, expected, **options):
    """Partial credit: 0.5 for all Bills emails starred, 0.5 for a star-action filter existing."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('all_starred'):
        score += 0.5
    if result.get('filter_has_star_action'):
        score += 0.5
    return score

def check_play_and_exit__ffd6f40d9680fbe09c9ab1f8d445d696_qw35sft2_3a30b889(actual_config_path, rule):
    """
    Validate VLC play-and-exit setting via vlcrc.
    play-and-exit=0: VLC stays open after playback (correct)
    play-and-exit=1: VLC auto-closes after playback (initial/wrong state)
    """
    try:
        with open(actual_config_path, 'rb') as f:
            config_text = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_2ef5dd.error('Failed to read VLC config: %s', e)
        return 0.0
    expected = rule.get('expected_play_and_exit', 0)
    if isinstance(expected, int):
        expected = str(expected)
    play_and_exit = '0'
    for line in config_text.split('\n'):
        stripped = line.strip()
        if stripped.startswith('#'):
            continue
        if 'play-and-exit=' in stripped:
            play_and_exit = stripped.split('=')[-1].strip()
    return 1.0 if play_and_exit == expected else 0.0

def check_global_keys_play_stop__0ab6bfa84809c927c91a6a1387eadc16_qw35sft2_1540c5d9(result, expected, **options):
    """
    Check if global hotkeys for BOTH Play/Pause AND Stop are set in the VLC config.
    Partial credit: 0.5 for each hotkey set.
    result: path to the vlcrc file
    expected: rules dict (unused, presence check only)
    """
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_e8a2da.error(f'Failed to read VLC config: {e}')
        return 0.0
    play_pause_set = False
    stop_set = False
    for line in config_file.split('\n'):
        if line.startswith('#') or '=' not in line:
            continue
        if 'global-key-play-pause=' in line:
            val = line.split('=')[-1].strip()
            if val:
                play_pause_set = True
        if 'global-key-stop=' in line:
            val = line.split('=')[-1].strip()
            if val:
                stop_set = True
    score = 0.0
    if play_pause_set:
        score += 0.5
    if stop_set:
        score += 0.5
    logger_qw35sft2_e8a2da.info(f'play_pause_set={play_pause_set}, stop_set={stop_set}, score={score}')
    return score

def check_global_key_next__d8953e567823b8ea130733a22ba09b75_qw35sft2_1b503325(result, expected, **options):
    """
    Check if a global hotkey for 'Next' track is set in the VLC config file.
    result: path to the vlcrc file (string)
    expected: rules dict with 'expected_global_key_next' (int 1 = set, 0 = not set)
    """
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_1e51cc.error(f'Failed to read VLC config: {e}')
        return 0.0
    global_key_next = '0'
    for line in config_file.split('\n'):
        if line.startswith('#') or '=' not in line:
            continue
        if 'global-key-next=' in line:
            val = line.split('=')[-1].strip()
            global_key_next = '0' if val == '' else '1'
    expected_value = expected.get('expected_global_key_next', 1)
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    logger_qw35sft2_1e51cc.info(f'global-key-next status: {global_key_next}, expected: {expected_value}')
    return 1.0 if global_key_next == expected_value else 0.0

def check_global_key_vol_up__d76c6ff727ecd69650cf5e0e8eb5e19a_qw35sft2_b23ce3cf(result, expected, **options):
    """
    Check if a global hotkey for 'Volume Up' is set in the VLC config file.
    result: path to the vlcrc file
    expected: rules dict with 'expected_global_key_vol_up' (int 1 = set, 0 = not set)
    """
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_9d1c6a.error(f'Failed to read VLC config: {e}')
        return 0.0
    global_key_vol_up = '0'
    for line in config_file.split('\n'):
        if line.startswith('#') or '=' not in line:
            continue
        if 'global-key-vol-up=' in line:
            val = line.split('=')[-1].strip()
            global_key_vol_up = '0' if val == '' else '1'
    expected_value = expected.get('expected_global_key_vol_up', 1)
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    logger_qw35sft2_9d1c6a.info(f'global-key-vol-up status: {global_key_vol_up}, expected: {expected_value}')
    return 1.0 if global_key_vol_up == expected_value else 0.0

def check_play_pause_and_recording__0ed5ebba8c137cdb9537f1717121259d_qw35sft2_4ed02195(result, expected, **options):
    """
    Check if global Play/Pause hotkey is set AND recording folder matches expected path.
    Partial credit: 0.5 for each condition satisfied.
    result: path to the vlcrc file
    expected: rules dict with optional 'expected_recording_path' (default '/home/user/Desktop')
    """
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_849761.error(f'Failed to read VLC config: {e}')
        return 0.0
    play_pause_set = False
    recording_correct = False
    expected_recording_path = expected.get('expected_recording_path', '/home/user/Desktop')
    for line in config_file.split('\n'):
        if line.startswith('#') or '=' not in line:
            continue
        if 'global-key-play-pause=' in line:
            val = line.split('=')[-1].strip()
            if val:
                play_pause_set = True
        if 'input-record-path=' in line:
            val = line.split('=')[-1].strip()
            if val == expected_recording_path:
                recording_correct = True
    score = 0.0
    if play_pause_set:
        score += 0.5
    if recording_correct:
        score += 0.5
    logger_qw35sft2_849761.info(f'play_pause_set={play_pause_set}, recording_correct={recording_correct}, score={score}')
    return score

def check_global_keys_play_next__2b7ee2d73c0169078c29f1690d78cdca_qw35sft2_1bb2ac84(result, expected, **options):
    """
    Check if global hotkeys for BOTH Play/Pause AND Next are set in the VLC config.
    Partial credit: 0.5 for each hotkey set.
    result: path to the vlcrc file
    expected: rules dict (unused, presence check only)
    """
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_08d9ad.error(f'Failed to read VLC config: {e}')
        return 0.0
    play_pause_set = False
    next_set = False
    for line in config_file.split('\n'):
        if line.startswith('#') or '=' not in line:
            continue
        if 'global-key-play-pause=' in line:
            val = line.split('=')[-1].strip()
            if val:
                play_pause_set = True
        if 'global-key-next=' in line:
            val = line.split('=')[-1].strip()
            if val:
                next_set = True
    score = 0.0
    if play_pause_set:
        score += 0.5
    if next_set:
        score += 0.5
    logger_qw35sft2_08d9ad.info(f'play_pause_set={play_pause_set}, next_set={next_set}, score={score}')
    return score

def check_ext_and_multi_settings__afc2fb3b68b53df8d476e05f07933aff_qw35sft2_4a8a2864(result, expected, **options):
    """Check extension (0.34) + tabSize (0.33) + formatOnSave (0.33)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'undefined_publisher.test')
    if ext_id in result.get('ext_list', ''):
        score += 0.34
    expected_tab = expected.get('tab_size')
    actual_tab = result.get('tab_size')
    if expected_tab is not None and actual_tab is not None:
        try:
            if int(actual_tab) == int(expected_tab):
                score += 0.33
        except (TypeError, ValueError):
            pass
    expected_fos = expected.get('format_on_save')
    actual_fos = result.get('format_on_save')
    if expected_fos is not None and actual_fos == expected_fos:
        score += 0.33
    return min(score, 1.0)

def check_dual_replace__ae0b3ab6740bff99e232fccee3b36aad_qw35sft2_358b4d9f(result, expected, **options):
    """Check that two text replacements were made in the file (partial credit 0.5 each)."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    forbidden1 = expected.get('forbidden1', '')
    required1 = expected.get('required1', '')
    min_count1 = expected.get('min_count1', 1)
    if forbidden1 and forbidden1 not in content and (content.count(required1) >= min_count1):
        score += 0.5
    forbidden2 = expected.get('forbidden2', '')
    required2 = expected.get('required2', '')
    min_count2 = expected.get('min_count2', 1)
    if forbidden2 and forbidden2 not in content and (content.count(required2) >= min_count2):
        score += 0.5
    return min(score, 1.0)

def check_dual_replace__004ad4b9f1b7c0d57d366b67195586e3_qw35sft2_0f6dfe81(result, expected, **options):
    """Check that two text replacements were made in the file (partial credit 0.5 each)."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    forbidden1 = expected.get('forbidden1', '')
    required1 = expected.get('required1', '')
    min_count1 = expected.get('min_count1', 1)
    if forbidden1 and forbidden1 not in content and (content.count(required1) >= min_count1):
        score += 0.5
    forbidden2 = expected.get('forbidden2', '')
    required2 = expected.get('required2', '')
    min_count2 = expected.get('min_count2', 1)
    if forbidden2 and forbidden2 not in content and (content.count(required2) >= min_count2):
        score += 0.5
    return min(score, 1.0)

def check_ext_and_wordwrap__324a83eafb9ff6ed07fafe7e199af4d5_qw35sft2_cd76830c(result, expected, **options):
    """Check autoDocstring installed (0.5) + editor.wordWrap matches expected value (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'njpwerner.autodocstring')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    actual_wrap = result.get('word_wrap')
    expected_wrap = expected.get('word_wrap')
    if expected_wrap is not None and actual_wrap == expected_wrap:
        score += 0.5
    return score

def check_randint_syntax__2c85ba49c59e818794c6aa64361a93f2_qw35sft2_68992a21(result, expected, **options):
    """Check that main.py contains the corrected np.random.randint call with proper comma-separated args."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else str(result)
    target = expected.get('expected_pattern', '')
    if not target:
        return 0.0
    return 1.0 if target in content else 0.0

def check_triple_replace__076bfd2d3ed2cd7e42a937a1dd80c8a9_qw35sft2_6ca0fcbf(result, expected, **options):
    """Check that three text replacements were made in the file (partial credit 0.33/0.34/0.33)."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    forbidden1 = expected.get('forbidden1', '')
    required1 = expected.get('required1', '')
    min_count1 = expected.get('min_count1', 1)
    if forbidden1 and forbidden1 not in content and (content.count(required1) >= min_count1):
        score += 0.33
    forbidden2 = expected.get('forbidden2', '')
    required2 = expected.get('required2', '')
    min_count2 = expected.get('min_count2', 1)
    if forbidden2 and forbidden2 not in content and (content.count(required2) >= min_count2):
        score += 0.34
    forbidden3 = expected.get('forbidden3', '')
    required3 = expected.get('required3', '')
    min_count3 = expected.get('min_count3', 1)
    if forbidden3 and forbidden3 not in content and (content.count(required3) >= min_count3):
        score += 0.33
    return min(score, 1.0)

def check_dual_ext__95ee9f441436911870f312520ed4f195_qw35sft2_fe2021e4(result, expected, **options):
    """Check two extensions installed (0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_list = result.get('ext_list', '')
    ext1 = expected.get('ext1', '')
    ext2 = expected.get('ext2', '')
    if ext1 and ext1 in ext_list:
        score += 0.5
    if ext2 and ext2 in ext_list:
        score += 0.5
    return score

def check_dual_replace__e470774862e9eb18ed4fb0ad458ab39f_qw35sft2_45510beb(result, expected, **options):
    """Check that two text replacements were made in the file (partial credit 0.5 each)."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    forbidden1 = expected.get('forbidden1', '')
    required1 = expected.get('required1', '')
    min_count1 = expected.get('min_count1', 1)
    if forbidden1 and forbidden1 not in content and (content.count(required1) >= min_count1):
        score += 0.5
    forbidden2 = expected.get('forbidden2', '')
    required2 = expected.get('required2', '')
    min_count2 = expected.get('min_count2', 1)
    if forbidden2 and forbidden2 not in content and (content.count(required2) >= min_count2):
        score += 0.5
    return min(score, 1.0)

def check_ext_and_wordwrap__6ac7db98d9670abbcf37d963dc27bc84_qw35sft2_dbef1ed3(result, expected, **options):
    """Check extension installed (0.5) + editor.wordWrap == 'on' (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'undefined_publisher.test')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    expected_wrap = expected.get('word_wrap', 'on')
    if result.get('word_wrap') == expected_wrap:
        score += 0.5
    return score

def check_mean_call__4c84b586b63f7c1cf9b6df39aedcd5ed_qw35sft2_c0f39136(result, expected, **options):
    """Check that main.py uses .mean(axis=1) (parentheses) instead of .mean[axis=1] (brackets)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else str(result)
    target = expected.get('expected_pattern', '')
    bad_pattern = expected.get('bad_pattern', '')
    if not content:
        return 0.0
    has_fix = target in content if target else True
    still_broken = bad_pattern in content if bad_pattern else False
    if has_fix and (not still_broken):
        return 1.0
    return 0.0

def check_dual_replace__d87d9b6cb6d6a7fc9d5d36ba0c96e0de_qw35sft2_3f330cae(result, expected, **options):
    """Check that two text replacements were made in the file (partial credit 0.5 each)."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    forbidden1 = expected.get('forbidden1', '')
    required1 = expected.get('required1', '')
    min_count1 = expected.get('min_count1', 1)
    if forbidden1 and forbidden1 not in content and (content.count(required1) >= min_count1):
        score += 0.5
    forbidden2 = expected.get('forbidden2', '')
    required2 = expected.get('required2', '')
    min_count2 = expected.get('min_count2', 1)
    if forbidden2 and forbidden2 not in content and (content.count(required2) >= min_count2):
        score += 0.5
    return min(score, 1.0)
