"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_font_color__175bcb8418fdbfeda65db92690d9b128', 'check_writer_highlight__6b6e5048ebb3170359c47d1ce863bb95', 'check_bold__1f6e651edc27326d3dac73dc9e29333b', 'check_flight_params__4c8f37f88d66f775179dafd23aadc29e', 'check_all_text_italic__bf5294aa1b721bbae785292729cb8020', 'check_text_bold__c4941aca377cb2eab0e421e15f393854', 'check_docx_district_list__5cd9f8324e43d00ab6826873ca723208', 'check_docx_reverse_merge__08f1db516ca13465963b37a9fdcaf66c', 'check_title_bold_center__c874bb9d2ad63d6767d6a39e179946a5', 'find_heading_font__6cbd3e1cf4cd04e2af7f29f623cbbd37', 'check_docx_footer_page_numbers__b8220411d07c53e298b6db0365205e12', 'check_writer_font_sizes__5b4d6572a569badb37aa2fd25bab4b6e', 'check_docx_multi_step__5ac00db9c05416b4087c6a9ef6d0dee7', 'check_default_font__fc8e11a362be23f66acd12dd6bff4956', 'check_docx_text_replace__3df1073706339202301b753c7942ef7f', 'check_docx_all_fonts__e754543a404f008975bc7bab115cbf09', 'check_docx_para_font__ec9fd979c9f4ede06619498c1e9b1202', 'check_header_italic__ee9130b48c4f13e968a7877f09e8d206', 'check_docx_gpt4_avg_table__75b704f601003cb29b05bac44a154402', 'check_text_bold__db0f57f6a04c8947ec0364b499605e29', 'check_docx_styled_merge__3faeea9d9342eb480cdf899557de9715', 'check_pdf_landscape_single_page__718b8e151527c1d36352c5376a93ead7', 'check_italic_subtitle__910b283898dbff47da6b9e5280c358bb', 'check_docx_first_para_alignment__58a0796c274d2d883482eaaa895cfb2f', 'check_writer_title_alignment__3a67dac16aed3d74bffb9cb95c614713', 'check_odt_conversion__016af9433419594b3dc790848ad51e3b', 'check_bg_font_color__d8eda1446480065d3fccc5f3a4fe966e', 'check_docx_heading_bold__0d3391a68f9d89ec2a40039cb5c06d2d', 'check_docx_default_font__450e6c1e4a27707dc54b625806c54889', 'check_italic_bold__7532d291294fb4f3bb8e25015109f986', 'check_writer_title_fontsize__22037d9b737ef0f8bf33646ac10692bd', 'check_writer_first_para_bold__974838092a7939215ca9c77d0b32d057', 'check_docx_header__2308ac81a92b82ba916e71c364d68457', 'check_docx_selective_merge__1d49dccb31c44fca7886c9c825e7148f', 'check_italic_font_name__f2fb15859784992fd6cf15bd171d4b91', 'check_docx_para_alignment__4ed04700852dd5a18fbdc223aac7b1c3', 'check_title_bold__94187b1698fde7f55883e20a5cb9822f', 'find_default_font_size__c661aabf93261d4afd20c3babf830aac', 'check_docx_to_pdf__8f355d6af58b179cf68cc37725d92b0c', 'check_bold_formatting__06366a1cfa1039e4258f94612f6055b9', 'check_content_bold__ca58c5a207cc4d7f5ab415f13aba66b0', 'check_docx_font_size__0625ebacf8174a330d7b47b181788578', 'check_writer_font__2a303ca90381d050b859a1a253c0648a', 'check_title_bold__675486c95c0c9441584aa662f56cd677', 'check_flight_params__c2e765147c6269817e33ad57cd24dd1a', 'check_docx_page_numbers__18ec7a61b2f5504dc6e64dce31cd1e00', 'check_title_bold_and_size__6146cd6c34c0bc769a02e73ec8ef6f70', 'check_docx_font__c5ed5054f5dbce73ddb0483f71ff6f1a', 'check_title_font_size__3af9d622838a82de0731abee0010dddc', 'check_title_underline__b1d7e7e1d48378504a47d3ccfb3cf1be', 'check_writer_default_font__0f58b322b0ad4bcc09c937490ff9bfeb', 'check_font_size__a1c0faa26bf29b13cb4b3dfa5122d14a', 'check_footer_has_page_numbers__3205af8e0dea88590c54f355855d7cbf', 'check_docx_page_orientation__7f5136b7f1dcbd0181f2d126bcac0caf', 'check_docx_title_alignment__6fdb2c432a1af94ddff3a56f751c6296', 'check_pdf_page_count__a8da1aa98de37fa62b8d321ad0589eaf', 'check_docx_title_text__ca26e505bd69174991058c87e74f2a98', 'check_first_para_font__53ba1671387b371fc528911be537db49', 'check_docx_contains__c53695cdc8b25b507dc6019962dcaf31', 'check_docx_title_alignment__0a4a5cad7031961d3e6b497ec2c282f4', 'check_docx_font_info__2042c7daa4f47bea2c4fb2c5e4661307', 'check_docx_all_bold__aef5992d33ef93f956afe1a194ef0bce', 'check_docx_heading_alignment__de2b396948b313eaa0067057e04e4e77', 'check_docx_line__31146efdc87a4defa3f97309b7095c90', 'check_docx_italic__d7547f75d107bc34f4a878c280fffbc1', 'check_italic__39e2309fd148a4ecc88950c3fc66ef3b', 'check_budget_page_state__b76001106213215f01706a09eab507db', 'check_docx_contains__72fc0161fa321ae4bf01e0be489f2375', 'check_docx_table_bold__b70d187abe0811a915df187857b97f14', 'check_docx_contains__896d76d765c6621bc8feff2ba14e6be9', 'check_docx_para_bold__7a1ffa51886994c5f3667d404d37d94d', 'check_docx_text_and_image__2aeb38ddfad7ab2f7bd9e6f63f6cc1bf', 'check_docx_default_font__536a5616674c51027b6cacf28aa7ecd9', 'check_docx_title__b98b6fc5ce8ebe003cf7eaeacc71d919', 'check_footer_page_numbers__806bfdf7c0ee49dcce08f91a2edad9e6', 'check_docx_alignment__ae9bdf892a1bf6b8ab5ef93251d291e0', 'check_title_size_content_bold_bg__80858354e7972b18d1e56bf385980263', 'check_bg_fontsize__c22d8b15c916d1ef420eb42daeb27bf0', 'check_heading2_fonts__3ee05560ea0a8eae187c5c1b3f30c5d3', 'check_text_underline__1563ed9d69890bfb1ca48bd048aeee4a', 'check_font_color_match__d0c681af3bd2af47798202cb5ebe1988', 'check_docx_para_alignment__099018697289346087fe050b42a1dd16', 'check_docx_model_names_table__257a7a045a00e35dcfe1e7d02a439bcc', 'check_docx_text_formatting__92a351c44e937610849c0af1deb80b89', 'check_column_header_and_values__23bfc830003e8f489b53a39808f56b59', 'check_doc_file_list__d6bd2cbacad8b9b51e1e935f00211e4d', 'check_docx_gpt4_specific_scores__9467a3bac5aae80d6f5894cff6144c37', 'check_font_color__3892513c1cab5055660dfeaa188121fb', 'check_page_numbers__9de26a1b58798be1fa2f4b6b3d9e9ec9', 'check_pdf_page_count__ddacf8dd85e44a11a4923ff07ce0abff', 'check_all_fonts_changed__f0d9875ec4c08351c2fb33918a5793ea', 'check_header_corrected__e0d3491cdcd2d9a19af59b5d2c376712', 'check_underline__935b1aaa61bb5f50d9e66c69d9b5829b', 'check_document_font__2ae82910da68002280a1510e1ca61e99', 'check_flight_params__15e6e06fadd423d4e6a279cca6f8914d', 'check_title_font_color__dcc119efe7774eb3168a0af998c68532', 'check_has_page_numbers__ad71c47cc02b102c764d5154ce0d73c0', 'check_font_size__e0ab347b041b8b1f40c9ca0e5bc35d39', 'check_italic_color__575ab439795c172e41189face7efb98a', 'check_title_font__7f1ba412564f3dc4322d685ce4a3274f', 'check_hint_line_spacing__7fbc318f6c5b72e9f18e5cf28491b9c0', 'check_writer_title_italic__088f05c0fd2c52a92bc3dc8f8cc6719c', 'check_title_font__1b8a8292a0fd2eb13bf16e945356170a', 'check_startup_and_font__aa9af3d49c1447b85e91f3d729dc4ea1_qw35sft2_1947001b', 'check_diamond_campground_page__c0084be8e992e4c5b6a50c34c3f5f437_qw35sft2_c0d65494', 'check_chrome_dnt_and_font__f0db2c12ec96d3956f7c7b7b8c915fc1_qw35sft2_7b230a8b', 'check_font_and_dnt__baedb9f4aec1e7501564ae9f37c335ac_qw35sft2_a1f236ee', 'check_header_bold__988aa3731e056dbe2f61a28b2135e9ae_qw35sft2_54bf3421', 'check_bold_header_and_pdf__86ea3ac794566dae35e2163b9c7d3c94_qw35sft2_0702a49e', 'check_bold_and_csv__8f6a3ccc5fb938975de38e8fbcc8f580_qw35sft2_69080047', 'check_employee_split_bold_header__fc565b338fbb3db18703aef2eaa623cb_qw35sft2_5a37c450', 'check_bold_header__771ce45281d75f3370a131f92752ad7a_qw35sft2_36c09f03', 'check_header_bold__636602487d408f3cbef02e77b45f21ff_qw35sft2_0ad6b944', 'check_title_underline__42b6fd5c06cd718a123cb3a7892d354e_qw35sft2_3c48d0f0', 'check_title_bold_and_size__d1408966fdb4fd77bea2fd3e21b1ee11_qw35sft2_04364771', 'check_strikethrough_and_italic__be4183a9c6982ba683532f14bc50d4bc_qw35sft2_4f1f568a', 'check_title_bold__78dfac329c5cd8b1c5efec5b79600fdd_qw35sft2_99d8b746', 'check_title_font_size__f083039072ffd98eec2d101bb28501a7_qw35sft2_4d31c951', 'check_writer_footer_body_spacing__e86d756354df667eec5758969f058362_qw35sft2_7eb35045', 'check_title_run_font__ac8fd7a53908049060d14f1a620d5f79_qw35sft2_f9a68cbe', 'check_para0_underline__0ee57f85e0ff64364a5561938bc94f89_qw35sft2_e008e549', 'check_docx_last_line_value__9ae69d0ef002958997c53c817727d757_qw35sft2_933992b3', 'check_docx_subscript_and_title_italic__5faa01d8a062ee366a1e56b39b81f8d9_qw35sft2_9e3c6cfb', 'check_table_7x5_with_headers__a4c465e58afc03255b4d0ddc46d5d525_qw35sft2_2ca96022', 'check_writer_basic_fonts__cb1a08b69b2e8e76cab6534a166b3ea6_qw35sft2_0eb16eec', 'check_docx_lower_and_page_nums__4155f473baa624a7c08f0169367bee57_qw35sft2_31960405', 'check_font_and_align__2b822025bba1f4240036b27d0339b871_qw35sft2_f2ba7288', 'check_docx_table_italic_col__838d41f5c432af14f1eb478fa956e99e_qw35sft2_7a2cc52c', 'check_writer_heading_body_align__d846f8a9d2cb6584815a6f17cfb8e80a_qw35sft2_481a2564', 'check_docx_first_para_strike__efe155f44671b1c09a26d5fefbd2fc44_qw35sft2_04b52bc6', 'check_all_italic_size_14__b8546f99d1a88f06fb0b6ffbfa13cf55_qw35sft2_ff244508', 'check_odt_highlight_italic__e4e69693f939d9cc32e39d11fb21f92a_qw35sft2_86960b9e', 'check_three_para_spacing__e4c36138929fcbc781e762a8996148c1_qw35sft2_f2c7fe41', 'check_writer_extended__d40d48ebd32c5cdd693a8c4ab4565d6c_qw35sft2_577878c3', 'check_docx_spacing_arial__d472f99c5ab0bbd719c99178119f321f_qw35sft2_1b4c6542', 'check_writer_break_and_notes__4e9ea3f60c13559abc5abe689f4b9ea2_qw35sft2_d173d01c', 'check_writer_default_and_list_fonts__12134317917b6b593b5731418bb79a41_qw35sft2_f9b4e005', 'check_docx_first_line_contains__d29d0e1e9fda0acc872853f3a63d2906_qw35sft2_6abfbae7', 'check_writer_footer_title_align__54873b40ed0174b0628f230c8ad47868_qw35sft2_7c373a94', 'check_docx_lower_and_bold_title__5844d45da47ec52046a23219bd5db770_qw35sft2_3c6668d1', 'check_docx_subscript_and_title_center__49f57bac6148724412fcfb3eb162539d_qw35sft2_fa414f8f', 'check_titlecase_and_italic__7a4ad09a1377c8f0eb798dbce87a7dce_qw35sft2_84c30304', 'check_pdf_and_odt_exists__c65906d6f690144adc013e7fcea2301e_qw35sft2_02517358', 'check_docx_table_structure__638d8a63e2ebfdf518128094d93e23f6_qw35sft2_f58f4a90', 'check_docx_has_image__ed4f0e4477e03c56fe74b28d9d6a3444_qw35sft2_5c2f253c', 'check_para0_word_replaced__469a4ad5eae55be7063b8ec35b77b37d_qw35sft2_a80ef772', 'check_docx_multi_para_strike__d8440cb16db3c65fff5a2b530ef73072_qw35sft2_d2ea5e29', 'check_italic_size_and_underline__dce073995a5927fc743181d7c02c0659_qw35sft2_c208bbd1', 'check_font_and_size__a567f4bdee19e51d5ed6a587d8d38f9d_qw35sft2_224f4149', 'check_writer_heading_align_size__a3c76448cfe24f131aa7bba8054b5d03_qw35sft2_4fa6f53a', 'check_odt_highlight_font__582984f0118e076c9ea08b5f36393d3b_qw35sft2_d4cd9443', 'check_writer_three_goals__c8e0b7dd7367f091fd322e812e9986d0_qw35sft2_2a25ad24', 'check_docx_spacing_fontsize__122cebb30dd4bbe8d6067021381568cf_qw35sft2_c1d5de85', 'check_writer_break_and_font__52445c32d012b75894f7e753d9bb73ed_qw35sft2_a130d84f', 'check_titlecase_and_underline__d22455594f3cb66658bea4571465f8f1_qw35sft2_41d40f3d', 'check_writer_footer_pageno__035b6f0ae922d6e2d1f24326ec60c904_qw35sft2_be318513', 'check_docx_text_all_upper__cbc878ff0e7a736898047bffbd6f85f6_qw35sft2_32945c7e', 'check_para0_bold__42cf2947548fcbcc72e5e44a50eb60dd_qw35sft2_2fb424d7', 'check_docx_train_removed__aa11b4f1fde2d6cc216fd8dac61371d5_qw35sft2_742684e6', 'check_writer_font_and_alignment__70647269807c892c9cd9454bb994336f_qw35sft2_b7538338', 'check_image_and_fontsize__285ee55215732bb31d8255226d83a755_qw35sft2_8e5c91c4', 'check_docx_table_header_align__83874f16b468982748e51b1480830a3a_qw35sft2_1e849fbb', 'check_docx_subscript_and_bold_heading__fa57f784af7171fa0e677f21369b4177_qw35sft2_2b7094fd', 'check_all_italic_size_16__e2de398ff1a46d21a89c93f18cb16653_qw35sft2_97c614d6', 'check_font_and_bold__3a9eb9ceb78a9dd18ebdb09636bdbe29_qw35sft2_c7731e94', 'check_writer_heading_align_font__fb136ac47c87dcc7e787c09564b245bc_qw35sft2_5267f139', 'check_docx_strike_bold_chain__30608d866c349bd16efe53a1c72a27d7_qw35sft2_dde7a53f', 'check_odt_highlight_strikethrough__a30ceeb68a8eeb0b8a512ab61a2387e1_qw35sft2_e68b1a1e', 'check_writer_ref_added__dab123b7868d167252ed777095add0ba_qw35sft2_57d735cd', 'check_docx_spacing_bold_intro__7ab11a38694adeb4332550a69fd5abfe_qw35sft2_646f94b4', 'check_writer_break_and_title_italic__0f7fdf2d11eca0ddff6ddca8ba1c7a92_qw35sft2_0d4a6598', 'check_docx_sentence_case__121c5b95c7fa376e725e2a74f1f18a20_qw35sft2_fe4e5dd5', 'check_writer_footer_title_italic__6a6cc374d796e945ca1edd575879589a_qw35sft2_1abc48e5', 'check_docx_last_line_value__99126397762866f43ef4e6109070d76d_qw35sft2_f21f6be2', 'check_docx_last_line_value__e6176eae3a47bbdd8f12a06b8b082ed4_qw35sft2_5a3ab949', 'check_writer_font_and_pagebreak__e2b27ec72c919c5a4bd4f151e7ff64b6_qw35sft2_99965063', 'check_docx_table_bold_header__d0985f126b123d58872633f0350725ab_qw35sft2_a8a693fa', 'check_docx_subscript_title_and_body__cb7672ac12071a9e4c398d11dc4c470e_qw35sft2_7e8df99d', 'check_font_and_italic__c3632663016de6f20637b59e83145d05_qw35sft2_44c1680c', 'check_doc_page_breaks__b7c367c074a6362d7d9c85a08867a367_qw35sft2_f869087d', 'check_italic_and_title_size__792e97b1e8b831bfb11064f6cde55b00_qw35sft2_dc097e8b', 'check_writer_heading_align_italic__edb463affe3b938af85226bb52cb1b03_qw35sft2_e5e99d63', 'check_para0_italic__d2318c984412daccf01ad6c479224170_qw35sft2_d6dcf621', 'check_docx_last_para_bold__c977b8be915d8296c568265b630cf188_qw35sft2_909f71e3', 'check_docx_spacing_italic_conclusion__bb617d210691b208d4c61100467118c1_qw35sft2_c1ca409b', 'check_writer_ref_footer__78615d30a3700f795271d9e8e63e0686_qw35sft2_4066d575', 'check_writer_page_break_count__e042d7b442613a05635c554b051c41b1_qw35sft2_d2800da9', 'check_odt_highlight_fontsize__4ce1d37833ab4f6aa409d7454a343799_qw35sft2_1f515809', 'check_line_spacing__80aefd2088f27d1be89724efb49ed091_qw35sft2_c6d4321c', 'check_writer_footer_and_header__5d6628525b72a5805637db6ac940e950_qw35sft2_b0dc9230', 'check_pdf_and_footer_pagenum__0bb862e1956085267ced241046059b10_qw35sft2_aca248a7', 'check_image_and_pagenumbers__6390ff8ee787412726544ea6ba43db0d_qw35sft2_ee2dbcca', 'check_docx_subscript_and_heading_underline__b75f28775a1b3fe6faeb633d14a05fab_qw35sft2_266881b5', 'check_titlecase_and_bold__2777e85b511814778d7406b21395647f_qw35sft2_af2391c1', 'check_footer_page_numbers__699ab5651848f548e06466de9777d875_qw35sft2_417f7c36', 'check_docx_dedup__f922eeeed3d49013fd1e13103dbfb120_qw35sft2_9b95ea64', 'check_writer_font_and_footer__d1ae06a2cb3f08c6662cf614ce58e285_qw35sft2_70843d4e', 'check_italic_size_and_bold__57bf430ce2e461a40bf234942a8ba4ad_qw35sft2_5635b7b0', 'check_font_and_size__5934b0de41c172c0d1662adc3bbc874f_qw35sft2_1bc37137', 'check_writer_heading_align_body_bold__e1c39faa270c25698b597d87598bba49_qw35sft2_9c96e15c', 'check_para0_font_size__7211ff5a62d6ff910c67632b31846d1c_qw35sft2_aa8614af', 'check_docx_table_last_row__e0b81f05c0001516085ceb6f32f7b48c_qw35sft2_bc759d94', 'check_docx_second_para_italic__aa396abb2d9562e18e97067a7c8c6fe9_qw35sft2_83fb48d7', 'check_docx_spacing_center_intro__5f00d786b29854e31074fa3cfeefea1b_qw35sft2_55dd34c9', 'check_writer_citation14__21b44c209dfc410a39a6375bc0042826_qw35sft2_2c6f58ae', 'check_docx_with_header__2354b786c1a1e08afc94d3d722cbb7a6_qw35sft2_f6f184cb', 'check_docx_title__40a84e6a55059151959a2b459d4c099e_qw35sft2_3441cc96', 'check_python_docs_large_font__492670c8dd0e3374d756c80fee290bb9_qw35sft2_4ecbb99e', 'check_docx_duration__48db15ca4dac03d9c5bae1d76e883852_qw35sft2_7772db27', 'check_docx_gemini_paragraphs__abfc7551de3c0f45f173130d54033329_qw35sft2_d5a836b4', 'check_pandoc_installed__79ba7b4d76657204de8c0df9dc21c02e_qw35sft2_b9982ce5', 'check_vscode_theme_font_wrap__12b042cae58f794626351d3c568391b7_qw35sft2_b1232fc0', 'check_ext_and_fontsize__fe14a817663aeade8a20cb0f5baad2b6_qw35sft2_9dd1b11f', 'check_workspace_and_fontsize__4cb0241f3e1c4176b65966dc96627af4_qw35sft2_c55e228e', 'check_vscode_debug_focus_font__802820e3fd86caf9b92a9751803125a2_qw35sft2_84c57be3', 'check_ext_and_fontsize__bb23bc65a82d76a52cba97bb0a6bf9bc_qw35sft2_439af9b1', 'check_vscode_exclude_fontsize__33cb4b35f94afe16fa1fe71880064b0b_qw35sft2_0eef60d3']

def check_font_color__175bcb8418fdbfeda65db92690d9b128(result, expected, **options):
    """Check if all text runs have the expected font color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    colors = result.get('colors', [])
    expected_color = expected.get('expected_color', '').upper()
    if not colors:
        return 0.0
    matching = sum((1 for c in colors if c and c.upper() == expected_color))
    return matching / len(colors)

def check_writer_highlight__6b6e5048ebb3170359c47d1ce863bb95(result, expected, **options):
    """Check if target text instances are highlighted in yellow.

    Scoring:
    - Proportional: (highlighted_target_count / target_total_count)
    - 1.0 if all instances of target text are highlighted
    - 0.0 if no instances are highlighted or error
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    target_total = result.get('target_total_count', 0)
    target_highlighted = result.get('target_highlighted_count', 0)
    if target_total == 0:
        return 0.0
    score = min(target_highlighted / target_total, 1.0)
    return score

def check_bold__1f6e651edc27326d3dac73dc9e29333b(result, expected, **options):
    """Check if all text runs in table are bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('all_bold', False):
        return 1.0
    total = result.get('total_runs', 0)
    bold = result.get('bold_runs', 0)
    if total == 0:
        return 0.0
    return bold / total

def check_flight_params__4c8f37f88d66f775179dafd23aadc29e(result, expected, **options):
    """Check flight search URL parameters with partial credit.
    Gives 0.5 for correct fromStation, 0.5 for correct toStation.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_from = expected.get('fromStation', '')
    expected_to = expected.get('toStation', '')
    if expected_from and result.get('fromStation', '') == expected_from:
        score += 0.5
    if expected_to and result.get('toStation', '') == expected_to:
        score += 0.5
    return score

def check_all_text_italic__bf5294aa1b721bbae785292729cb8020(result, expected, **options):
    """Check if all text in the document is italic."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    ratio = result.get('ratio', 0.0)
    total = result.get('total_runs', 0)
    if total == 0:
        return 0.0
    if ratio >= 0.95:
        return 1.0
    elif ratio >= 0.7:
        return 0.5
    return 0.0

def check_text_bold__c4941aca377cb2eab0e421e15f393854(result, expected, **options):
    """Check if text runs are all bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('all_bold', False)
    return 1.0 if actual_bold == expected_bold else 0.0

def check_docx_district_list__5cd9f8324e43d00ab6826873ca723208(result, expected, **options):
    """Check if the docx contains all expected district names with partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_items = expected.get('expected_items', [])
    full_text = result.get('text', '')
    if not full_text or not expected_items:
        return 0.0
    matched = 0
    needs_fuzzy = []
    for item in expected_items:
        if item in full_text:
            matched += 1
        else:
            needs_fuzzy.append(item)
    if needs_fuzzy:
        from rapidfuzz import fuzz
        for item in needs_fuzzy:
            best_score = 0
            for para in result.get('paragraphs', []):
                score = fuzz.partial_ratio(item, para) / 100.0
                best_score = max(best_score, score)
            if best_score >= 0.8:
                matched += 1
    return matched / len(expected_items) if expected_items else 0.0

def check_docx_reverse_merge__08f1db516ca13465963b37a9fdcaf66c(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    first_150 = result.get('first_150', '').lower()
    last_150 = result.get('last_150', '').lower()
    first_expected = expected.get('first_snippet', '').lower()
    last_expected = expected.get('last_snippet', '').lower()
    if first_expected and first_expected in first_150:
        score += 0.2
    if last_expected and last_expected in last_150:
        score += 0.2
    text = result.get('text', '').lower()
    required_snippets = expected.get('required_snippets', [])
    if required_snippets:
        found = sum((1 for s in required_snippets if s.lower() in text))
        score += found / len(required_snippets) * 0.3
    expected_font_size = expected.get('expected_font_size')
    if expected_font_size is not None:
        font_sizes = result.get('font_sizes', [])
        if font_sizes and all((abs(fs - expected_font_size) < 0.5 for fs in font_sizes)):
            score += 0.3
    return min(score, 1.0)

def check_title_bold_center__c874bb9d2ad63d6767d6a39e179946a5(result, expected, **options):
    """Check if the title is bold and center-aligned. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_bold'):
        score += 0.5
    if result.get('is_centered'):
        score += 0.5
    return score

def find_heading_font__6cbd3e1cf4cd04e2af7f29f623cbbd37(config_file_path, rules):
    """Find the default heading font in LibreOffice Writer."""
    heading_font = None
    expected_font = rules['font_name']
    if not config_file_path:
        return 0
    try:
        tree = ET.parse(config_file_path)
        root = tree.getroot()
        namespace = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('./item[@oor:path="/org.openoffice.Office.Writer/DefaultFont/Standard"]', namespace):
            for prop in elem.findall('./prop[@oor:name="sHeading"]', namespace):
                for value in prop.findall('value', namespace):
                    heading_font = value.text
    except Exception as e:
        logger.error(f'Error: {e}')
    return 1 if heading_font == expected_font else 0

def check_docx_footer_page_numbers__b8220411d07c53e298b6db0365205e12(result, expected, **options):
    """Check if the document has page numbers in the footer."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if isinstance(result, dict) and result.get('has_page_number'):
        return 1.0
    return 0.0

def check_writer_font_sizes__5b4d6572a569badb37aa2fd25bab4b6e(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_size', 12)
    sizes = result.get('run_sizes', [])
    style_size = result.get('style_size')
    if not sizes and style_size is None:
        return 0.0
    if sizes:
        matching = sum((1 for s in sizes if abs(s - expected_size) < 0.5))
        return matching / len(sizes)
    if style_size is not None and abs(style_size - expected_size) < 0.5:
        return 1.0
    return 0.0

def check_docx_multi_step__5ac00db9c05416b4087c6a9ef6d0dee7(result, expected, **options):
    """Check docx contains both a header and address items with partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    full_text = result.get('text', '')
    if not full_text:
        return 0.0
    score = 0.0
    needs_fuzzy_header = False
    needs_fuzzy_items = []
    expected_header = expected.get('expected_header', '')
    if expected_header:
        if expected_header in full_text:
            score += 0.3
        else:
            needs_fuzzy_header = True
    expected_items = expected.get('expected_items', [])
    if expected_items:
        item_weight = 0.7 / len(expected_items)
        for item in expected_items:
            if item in full_text:
                score += item_weight
            else:
                needs_fuzzy_items.append((item, item_weight))
    if needs_fuzzy_header or needs_fuzzy_items:
        from rapidfuzz import fuzz
        if needs_fuzzy_header and expected_header:
            fuzzy_score = fuzz.partial_ratio(expected_header, full_text) / 100.0
            if fuzzy_score >= 0.8:
                score += 0.3 * fuzzy_score
        for (item, weight) in needs_fuzzy_items:
            best = 0
            for para in result.get('paragraphs', []):
                s = fuzz.partial_ratio(item, para) / 100.0
                best = max(best, s)
            if best >= 0.8:
                score += weight * best
    return min(score, 1.0)

def check_default_font__fc8e11a362be23f66acd12dd6bff4956(result, expected, **options):
    """Check if the default font matches the expected font name.

    Scoring: 1.0 if font matches (case-insensitive), 0.0 otherwise.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_font = result.get('default_font', '')
    expected_font = expected.get('expected_font', '')
    if not actual_font or not expected_font:
        return 0.0
    if actual_font.lower().strip() == expected_font.lower().strip():
        return 1.0
    return 0.0

def check_docx_text_replace__3df1073706339202301b753c7942ef7f(result, expected, **options):
    """Check if text replacement was done correctly.

    Scoring:
    - 0.5: '<add here>' is removed
    - 0.5: '(3)' is present in the Pennington paragraph
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_add_here', True):
        score += 0.5
    pennington = result.get('pennington_context', '')
    expected_text = expected.get('replacement_text', '(3)')
    if expected_text in pennington:
        score += 0.5
    return min(score, 1.0)

def check_docx_all_fonts__e754543a404f008975bc7bab115cbf09(result, expected, **options):
    """Check if all fonts in the document match the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', 'Arial')
    fonts = result.get('fonts', [])
    if not fonts:
        return 0.0
    if len(fonts) == 1 and fonts[0] == expected_font:
        return 1.0
    if expected_font in fonts:
        return 0.5
    return 0.0

def check_docx_para_font__ec9fd979c9f4ede06619498c1e9b1202(result, expected, **options):
    """Check if paragraph font matches expected font name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', '')
    total_chars = result.get('total_chars', 0)
    if total_chars == 0:
        return 0.0
    font_counts = result.get('font_counts', {})
    matched_chars = font_counts.get(expected_font, 0)
    ratio = matched_chars / total_chars
    if ratio >= 0.9:
        return 1.0
    elif ratio > 0.0:
        return ratio
    return 0.0

def check_header_italic__ee9130b48c4f13e968a7877f09e8d206(result, expected, **options):
    """Check if header cells are italic. Partial credit per cell."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    italic_states = result.get('italic_states', {})
    expected_cells = expected.get('expected_italic_cells', [])
    if not expected_cells:
        return 0.0
    correct = 0
    for cell_ref in expected_cells:
        if italic_states.get(cell_ref) is True:
            correct += 1
    return correct / len(expected_cells)

def check_docx_gpt4_avg_table__75b704f601003cb29b05bac44a154402(result, expected, **options):
    """Check if docx has a table containing GPT-4 name and its average score.
    Partial credit: 0.5 for table with Gpt-4, 1.0 for table with Gpt-4 and avg score.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    tables = result.get('tables', [])
    if not tables:
        return 0.0
    expected_name = expected.get('model_name', 'Gpt-4').lower()
    expected_avg = expected.get('avg_score', '0.1557')
    score = 0.0
    found_name = False
    found_avg = False
    for table in tables:
        for row in table:
            row_lower = [cell.lower() for cell in row]
            row_joined = ' '.join(row_lower)
            if expected_name.lower() in row_joined:
                found_name = True
            for cell in row:
                cell_clean = cell.strip()
                try:
                    val = float(cell_clean)
                    expected_val = float(expected_avg)
                    if abs(val - expected_val) < 0.01:
                        found_avg = True
                except (ValueError, TypeError):
                    pass
    if found_name:
        score += 0.5
    if found_avg:
        score += 0.5
    return min(score, 1.0)

def check_text_bold__db0f57f6a04c8947ec0364b499605e29(result, expected, **options):
    """Check if all text runs have the expected bold status."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    bold_values = result.get('bold_values', [])
    expected_bold = expected.get('expected_bold', True)
    if not bold_values:
        return 0.0
    matching = sum((1 for b in bold_values if b == expected_bold))
    return matching / len(bold_values)

def check_docx_styled_merge__3faeea9d9342eb480cdf899557de9715(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    text = result.get('text', '').lower()
    required_snippets = expected.get('required_snippets', [])
    if required_snippets:
        found = sum((1 for s in required_snippets if s.lower() in text))
        score += found / len(required_snippets) * 0.34
    expected_font_size = expected.get('expected_font_size')
    if expected_font_size is not None:
        font_sizes = result.get('font_sizes', [])
        if font_sizes and all((abs(fs - expected_font_size) < 0.5 for fs in font_sizes)):
            score += 0.33
    expected_font_name = expected.get('expected_font_name', '').lower()
    if expected_font_name:
        font_names = [fn.lower() for fn in result.get('font_names', [])]
        if font_names and any((expected_font_name in fn for fn in font_names)):
            score += 0.33
    return min(score, 1.0)

def check_pdf_landscape_single_page__718b8e151527c1d36352c5376a93ead7(result, expected, **options):
    """Check PDF is single page and landscape orientation. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_pages = expected.get('expected_pages', 1)
    expected_landscape = expected.get('expected_landscape', True)
    if result.get('page_count') == expected_pages:
        score += 0.5
    if result.get('is_landscape') == expected_landscape:
        score += 0.5
    return min(score, 1.0)

def check_italic_subtitle__910b283898dbff47da6b9e5280c358bb(result, expected, **options):
    """Check if title is italic and subtitle matches expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_italic') is True:
        score += 0.5
    expected_subtitle = expected.get('expected_subtitle', '').strip().lower()
    actual_subtitle = (result.get('subtitle_text') or '').strip().lower()
    if expected_subtitle and expected_subtitle in actual_subtitle:
        score += 0.5
    return min(score, 1.0)

def check_docx_first_para_alignment__58a0796c274d2d883482eaaa895cfb2f(result, expected, **options):
    """Check if the first paragraph has the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'center')
    actual_alignment = result.get('alignment', '')
    if actual_alignment.lower() == expected_alignment.lower():
        return 1.0
    return 0.0

def check_writer_title_alignment__3a67dac16aed3d74bffb9cb95c614713(result, expected, **options):
    """Check if the title paragraph has the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'center')
    actual_alignment = result.get('alignment', '')
    if actual_alignment == expected_alignment:
        return 1.0
    return 0.0

def check_odt_conversion__016af9433419594b3dc790848ad51e3b(result, expected, **options):
    """Check that .doc files were converted to .odt format via command line."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_odt_count', 12)
    if result.get('command_found'):
        score += 0.4
    actual_count = result.get('odt_count', 0)
    if actual_count >= expected_count:
        score += 0.6
    elif actual_count > 0:
        score += 0.6 * (actual_count / expected_count)
    return min(score, 1.0)

def check_bg_font_color__d8eda1446480065d3fccc5f3a4fe966e(result, expected, **options):
    """Check background color and font color with tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    tolerance = expected.get('tolerance', 50)
    expected_bg = expected.get('expected_bg_color', '')
    actual_bg = result.get('bg_color', '')
    bg_dist = 999
    if actual_bg and expected_bg:
        try:
            (r1, g1, b1) = (int(actual_bg[0:2], 16), int(actual_bg[2:4], 16), int(actual_bg[4:6], 16))
            (r2, g2, b2) = (int(expected_bg[0:2], 16), int(expected_bg[2:4], 16), int(expected_bg[4:6], 16))
            bg_dist = ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
        except Exception:
            bg_dist = 999
    if bg_dist <= tolerance:
        score += 0.5
    expected_font = expected.get('expected_font_color', '')
    actual_font = result.get('font_color', '')
    font_dist = 999
    if actual_font and expected_font:
        try:
            (r1, g1, b1) = (int(actual_font[0:2], 16), int(actual_font[2:4], 16), int(actual_font[4:6], 16))
            (r2, g2, b2) = (int(expected_font[0:2], 16), int(expected_font[2:4], 16), int(expected_font[4:6], 16))
            font_dist = ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
        except Exception:
            font_dist = 999
    if font_dist <= tolerance:
        score += 0.5
    return min(score, 1.0)

def check_docx_heading_bold__0d3391a68f9d89ec2a40039cb5c06d2d(result, expected, **options):
    """Check if the References heading is bold.

    Scoring:
    - 1.0: heading is fully bold
    - 0.5: heading is partially bold (some runs bold)
    - 0.0: not bold at all or error
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    if result.get('is_bold') == expected_bold:
        return 1.0
    elif result.get('any_bold') and expected_bold:
        return 0.5
    return 0.0

def check_docx_default_font__450e6c1e4a27707dc54b625806c54889(result, expected, **options):
    """Check if document default font matches expected. Partial credit: 0.5 for style, 0.5 for body text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', '')
    score = 0.0
    if result.get('default_font') and result['default_font'].lower() == expected_font.lower():
        score += 0.5
    if result.get('most_common_font') and result['most_common_font'].lower() == expected_font.lower():
        score += 0.5
    return min(score, 1.0)

def check_italic_bold__7532d291294fb4f3bb8e25015109f986(result, expected, **options):
    """Check if all italic text also has bold formatting."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    italic_bold = result.get('italic_bold', [])
    if not italic_bold:
        return 0.0
    matching = sum((1 for b in italic_bold if b is True))
    return matching / len(italic_bold)

def check_writer_title_fontsize__22037d9b737ef0f8bf33646ac10692bd(result, expected, **options):
    """Check if the title font size matches the expected value.

    Scoring:
    - 1.0 if font size matches exactly
    - 0.0 if no match or error
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_fontsize', '18pt')
    actual_size = result.get('title_fontsize')
    if actual_size is None:
        return 0.0
    actual_norm = str(actual_size).strip().lower().replace(' ', '')
    expected_norm = str(expected_size).strip().lower().replace(' ', '')
    if actual_norm == expected_norm:
        return 1.0
    try:
        actual_num = float(actual_norm.replace('pt', '').replace('px', ''))
        expected_num = float(expected_norm.replace('pt', '').replace('px', ''))
        if abs(actual_num - expected_num) < 0.1:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_writer_first_para_bold__974838092a7939215ca9c77d0b32d057(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    runs_bold = result.get('runs_bold', [])
    if not runs_bold:
        return 0.0
    matching = sum((1 for b in runs_bold if b))
    return matching / len(runs_bold)

def check_docx_header__2308ac81a92b82ba916e71c364d68457(result, expected, **options):
    """Check if the document contains the expected header text in the first paragraph."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_text = expected.get('expected_header', '').lower().strip()
    first_para = result.get('first_paragraph', '').lower().strip()
    all_paragraphs = result.get('all_paragraphs', [])
    if expected_text and expected_text in first_para:
        return 1.0
    for para in all_paragraphs:
        if expected_text in para.lower().strip():
            return 0.5
    return 0.0

def check_docx_selective_merge__1d49dccb31c44fca7886c9c825e7148f(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    text = result.get('text', '').lower()
    required_snippets = expected.get('required_snippets', [])
    if required_snippets:
        found = sum((1 for s in required_snippets if s.lower() in text))
        score += found / len(required_snippets) * 0.4
    forbidden_snippets = expected.get('forbidden_snippets', [])
    if forbidden_snippets:
        absent = sum((1 for s in forbidden_snippets if s.lower() not in text))
        score += absent / len(forbidden_snippets) * 0.3
    else:
        score += 0.3
    expected_font_size = expected.get('expected_font_size')
    if expected_font_size is not None:
        font_sizes = result.get('font_sizes', [])
        if font_sizes and all((abs(fs - expected_font_size) < 0.5 for fs in font_sizes)):
            score += 0.3
    return min(score, 1.0)

def check_italic_font_name__f2fb15859784992fd6cf15bd171d4b91(result, expected, **options):
    """Check if italic text uses expected font and non-italic text uses expected font.
    Partial scoring: 0.5 for italic font match + 0.5 for non-italic font match.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    italic_fonts = result.get('italic_fonts', [])
    expected_italic_font = expected.get('expected_font', 'Arial')
    if italic_fonts:
        matching = sum((1 for f in italic_fonts if f == expected_italic_font))
        score += 0.5 * (matching / len(italic_fonts))
    non_italic_fonts = result.get('non_italic_fonts', [])
    expected_non_italic_font = expected.get('expected_non_italic_font', 'Times New Roman')
    if non_italic_fonts:
        matching = sum((1 for f in non_italic_fonts if f == expected_non_italic_font))
        score += 0.5 * (matching / len(non_italic_fonts))
    return min(score, 1.0)

def check_docx_para_alignment__4ed04700852dd5a18fbdc223aac7b1c3(result, expected, **options):
    """Check if paragraph alignment matches expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'CENTER')
    actual_alignment = result.get('alignment', '')
    if actual_alignment.upper() == expected_alignment.upper():
        return 1.0
    return 0.0

def check_title_bold__94187b1698fde7f55883e20a5cb9822f(result, expected, **options):
    """Check if the title paragraph is bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('all_bold', False)
    if actual_bold == expected_bold:
        return 1.0
    return 0.0

def find_default_font_size__c661aabf93261d4afd20c3babf830aac(config_file_path, rules):
    """Find the default font size in LibreOffice Writer.

    LibreOffice stores font sizes in registrymodifications.xcu in hundredths
    of a point (e.g. 14pt -> 1400). The metric converts the stored value
    before comparing against the expected size in points.
    """
    font_size = None
    expected_size = rules['font_size']
    if not config_file_path:
        return 0
    try:
        tree = ET.parse(config_file_path)
        root = tree.getroot()
        namespace = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('./item[@oor:path="/org.openoffice.Office.Writer/DefaultFont/Standard"]', namespace):
            for prop in elem.findall('./prop[@oor:name="Height"]', namespace):
                for value in prop.findall('value'):
                    try:
                        font_size = int(value.text)
                    except (ValueError, TypeError):
                        font_size = value.text
    except Exception as e:
        logger.error(f'Error: {e}')
    if font_size is not None and expected_size is not None:
        try:
            font_size_pt = int(font_size) / 100
            return 1 if font_size_pt == int(expected_size) else 0
        except (ValueError, TypeError):
            pass
    return 0

def check_docx_to_pdf__8f355d6af58b179cf68cc37725d92b0c(result, expected, **options):
    """Check that .docx files were converted to PDF format via command line."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_pdf_count', 3)
    if result.get('command_found'):
        score += 0.4
    actual_count = result.get('docx_pdf_count', 0)
    if actual_count >= expected_count:
        score += 0.6
    elif actual_count > 0:
        score += 0.6 * (actual_count / expected_count)
    return min(score, 1.0)

def check_bold_formatting__06366a1cfa1039e4258f94612f6055b9(result, expected, **options):
    """Check if text has bold formatting applied."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('all_bold', False)
    if actual_bold == expected_bold:
        return 1.0
    return 0.0

def check_content_bold__ca58c5a207cc4d7f5ab415f13aba66b0(result, expected, **options):
    """Check content text and bold formatting. Partial credit: 0.5 text, 0.5 bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    expected_bold = expected.get('expected_bold', True)
    actual_text = result.get('content_text')
    actual_bold = result.get('is_bold')
    if actual_text and expected_text:
        if actual_text.strip().lower() == expected_text.strip().lower():
            score += 0.5
    if actual_bold is not None and actual_bold == expected_bold:
        score += 0.5
    return min(score, 1.0)

def check_docx_font_size__0625ebacf8174a330d7b47b181788578(result, expected, **options):
    """Check if the first paragraph font size matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_size = result.get('font_size')
    expected_size = expected.get('expected_font_size')
    if actual_size is None or expected_size is None:
        return 0.0
    return 1.0 if abs(float(actual_size) - float(expected_size)) < 0.5 else 0.0

def check_writer_font__2a303ca90381d050b859a1a253c0648a(result, expected, **options):
    """Check if document font has been changed to the expected font.

    Scoring:
    - 1.0 if the expected font is found in used_fonts (partial: font was applied to some text)
    - 0.5 if font is default but other fonts still exist
    - 0.0 if expected font not found at all
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', 'Arial')
    used_fonts = result.get('used_fonts', [])
    default_font = result.get('default_font')
    single_font = result.get('single_font')
    if single_font and single_font.lower() == expected_font.lower():
        return 1.0
    font_found = any((f.lower() == expected_font.lower() for f in used_fonts))
    default_matches = default_font and default_font.lower() == expected_font.lower()
    if font_found and default_matches:
        return 0.8
    elif font_found:
        return 0.5
    elif default_matches:
        return 0.5
    return 0.0

def check_title_bold__675486c95c0c9441584aa662f56cd677(result, expected, **options):
    """Check if the title text is bold."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    all_bold = result.get('all_bold', False)
    if all_bold == expected_bold:
        return 1.0
    bold_states = result.get('bold_states', [])
    if not bold_states:
        return 0.0
    if expected_bold:
        return sum((1 for s in bold_states if s)) / len(bold_states)
    else:
        return sum((1 for s in bold_states if not s)) / len(bold_states)

def check_flight_params__c2e765147c6269817e33ad57cd24dd1a(result, expected, **options):
    """Check flight search URL parameters with partial credit.
    Gives 0.5 for correct fromStation, 0.5 for correct toStation.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_from = expected.get('fromStation', '')
    expected_to = expected.get('toStation', '')
    if expected_from and result.get('fromStation', '') == expected_from:
        score += 0.5
    if expected_to and result.get('toStation', '') == expected_to:
        score += 0.5
    return score

def check_docx_page_numbers__18ec7a61b2f5504dc6e64dce31cd1e00(result, expected, **options):
    """Check if document has page numbers in footer."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    has_page_numbers = result.get('has_page_numbers', False)
    expected_value = expected.get('has_page_numbers', True)
    return 1.0 if has_page_numbers == expected_value else 0.0

def check_title_bold_and_size__6146cd6c34c0bc769a02e73ec8ef6f70(result, expected, **options):
    """Check if title is bold and has expected font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_size = expected.get('expected_size', 16)
    if result.get('is_bold'):
        score += 0.5
    font_sizes = result.get('font_sizes', [])
    if font_sizes:
        avg_size = sum(font_sizes) / len(font_sizes)
        if abs(avg_size - expected_size) < 0.5:
            score += 0.5
    elif not font_sizes and result.get('is_bold'):
        pass
    return min(score, 1.0)

def check_docx_font__c5ed5054f5dbce73ddb0483f71ff6f1a(result, expected, **options):
    """Check if document font matches expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', 'Arial')
    default_font = result.get('default_font', '')
    if default_font and default_font.lower() == expected_font.lower():
        return 1.0
    run_fonts = result.get('run_fonts', {})
    if not run_fonts:
        return 0.0
    total = sum(run_fonts.values())
    matching = 0
    for (font, count) in run_fonts.items():
        if font.lower() == expected_font.lower():
            matching += count
    if total > 0:
        ratio = matching / total
        if ratio >= 0.8:
            return 1.0
        return ratio
    return 0.0

def check_title_font_size__3af9d622838a82de0731abee0010dddc(result, expected, **options):
    """Check title text and font size. Partial credit: 0.5 for text, 0.5 for font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    expected_size = expected.get('expected_font_size_pt')
    actual_text = result.get('title_text')
    actual_size = result.get('font_size_pt')
    if actual_text and expected_text:
        if actual_text.strip().lower() == expected_text.strip().lower():
            score += 0.5
    if actual_size is not None and expected_size is not None:
        if abs(float(actual_size) - float(expected_size)) <= 0.5:
            score += 0.5
    return min(score, 1.0)

def check_title_underline__b1d7e7e1d48378504a47d3ccfb3cf1be(result, expected, **options):
    """Check if the title text is underlined."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_underlined = expected.get('expected_underlined', True)
    all_underlined = result.get('all_underlined', False)
    if all_underlined == expected_underlined:
        return 1.0
    underline_states = result.get('underline_states', [])
    if not underline_states:
        return 0.0
    if expected_underlined:
        return sum((1 for s in underline_states if s)) / len(underline_states)
    else:
        return sum((1 for s in underline_states if not s)) / len(underline_states)

def check_writer_default_font__0f58b322b0ad4bcc09c937490ff9bfeb(result, expected, **options):
    """Check if body text font matches expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', '')
    actual_font = result.get('font', '')
    if not actual_font:
        return 0.0
    if actual_font.lower() == expected_font.lower():
        return 1.0
    return 0.0

def check_font_size__a1c0faa26bf29b13cb4b3dfa5122d14a(result, expected, **options):
    """Check if font sizes match expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_size')
    font_sizes = result.get('font_sizes', [])
    if not font_sizes:
        return 0.0
    matching = sum((1 for s in font_sizes if abs(s - expected_size) < 0.5))
    return matching / len(font_sizes)

def check_footer_has_page_numbers__3205af8e0dea88590c54f355855d7cbf(result, expected, **options):
    """Check if footer contains page numbers."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if result.get('has_page_number') is True:
        return 1.0
    return 0.0

def check_docx_page_orientation__7f5136b7f1dcbd0181f2d126bcac0caf(result, expected, **options):
    """Check if document page orientation matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_landscape = expected.get('all_landscape', True)
    actual = result.get('all_landscape', False)
    if actual == expected_landscape:
        return 1.0
    if expected_landscape and result.get('any_landscape', False):
        return 0.5
    return 0.0

def check_docx_title_alignment__6fdb2c432a1af94ddff3a56f751c6296(result, expected, **options):
    """Check if the title paragraph is center-aligned."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if isinstance(result, dict) and result.get('is_centered'):
        return 1.0
    return 0.0

def check_pdf_page_count__a8da1aa98de37fa62b8d321ad0589eaf(result, expected, **options):
    """Check if PDF exists and has the expected number of pages."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('file_exists'):
        score += 0.5
    else:
        return 0.0
    expected_pages = expected.get('expected_pages', 5)
    actual_pages = result.get('num_pages', 0)
    if actual_pages == expected_pages:
        score += 0.5
    return min(score, 1.0)

def check_docx_title_text__ca26e505bd69174991058c87e74f2a98(result, expected, **options):
    """Check if the document title matches the expected title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if not expected_title:
        return 0.0
    if actual_title == expected_title:
        return 1.0
    if actual_title.lower() == expected_title.lower():
        return 0.8
    if expected_title.lower() in actual_title.lower():
        return 0.5
    return 0.0

def check_first_para_font__53ba1671387b371fc528911be537db49(result, expected, **options):
    """Check if ALL runs in the first paragraph use the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    fonts = result.get('fonts', [])
    expected_font = expected.get('expected_font', 'Arial')
    if not fonts:
        return 0.0
    for font in fonts:
        if font is None or font.lower() != expected_font.lower():
            return 0.0
    return 1.0

def check_docx_contains__c53695cdc8b25b507dc6019962dcaf31(result, expected, **options):
    """Check if the docx text contains the expected book title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    text = result.get('text', '').strip().lower()
    expected_title = expected.get('expected_title', '').strip().lower()
    if not expected_title:
        return 0.0
    if expected_title in text:
        return 1.0
    words = expected_title.split()
    if not words:
        return 0.0
    matched = sum((1 for w in words if w in text))
    return matched / len(words) * 0.5

def check_docx_title_alignment__0a4a5cad7031961d3e6b497ec2c282f4(result, expected, **options):
    """Check if the title alignment matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'center')
    actual_alignment = result.get('alignment', '')
    if actual_alignment == expected_alignment:
        return 1.0
    return 0.0

def check_docx_font_info__2042c7daa4f47bea2c4fb2c5e4661307(result, expected, **options):
    """Check if all text uses the expected font and size. Partial credit: 0.5 for font, 0.5 for size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', '')
    expected_size = expected.get('expected_size_pt', 0)
    actual_fonts = result.get('fonts', [])
    actual_sizes = result.get('sizes_pt', [])
    if actual_fonts and len(actual_fonts) == 1 and (actual_fonts[0] == expected_font):
        score += 0.5
    if actual_sizes and len(actual_sizes) == 1 and (abs(actual_sizes[0] - expected_size) < 0.5):
        score += 0.5
    return min(score, 1.0)

def check_docx_all_bold__aef5992d33ef93f956afe1a194ef0bce(result, expected, **options):
    """Check if all text in the document is bold. Supports partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    total_runs = result.get('total_runs', 0)
    bold_runs = result.get('bold_runs', 0)
    if total_runs == 0:
        return 0.0
    return bold_runs / total_runs

def check_docx_heading_alignment__de2b396948b313eaa0067057e04e4e77(result, expected, **options):
    """Check if the heading has the expected alignment.

    Scoring:
    - 1.0: alignment matches expected
    - 0.0: otherwise
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_alignment = expected.get('expected_alignment', 'CENTER')
    actual_alignment = result.get('alignment', '')
    if actual_alignment.upper() == expected_alignment.upper():
        return 1.0
    return 0.0

def check_docx_line__31146efdc87a4defa3f97309b7095c90(result, expected, **options):
    """Check if the matched line matches expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('matched_line', '')
    if actual is None:
        return 0.0
    expected_text = expected.get('expected_text', '')
    actual_clean = actual.strip().lower()
    expected_clean = expected_text.strip().lower()
    if actual_clean == expected_clean:
        return 1.0
    return 0.0

def check_docx_italic__d7547f75d107bc34f4a878c280fffbc1(result, expected, **options):
    """Check if the first paragraph italic status matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_italic = result.get('is_italic', False)
    expected_italic = expected.get('expected_italic', True)
    return 1.0 if actual_italic == expected_italic else 0.0

def check_italic__39e2309fd148a4ecc88950c3fc66ef3b(result, expected, **options):
    """Check if all text runs are italic."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('all_italic', False):
        return 1.0
    total = result.get('total_runs', 0)
    italic = result.get('italic_runs', 0)
    if total == 0:
        return 0.0
    return italic / total

def check_budget_page_state__b76001106213215f01706a09eab507db(result, expected, **options):
    """Check budget.com page state with partial credit for URL and sort."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    url = result.get('url', '')
    expected_url_pattern = expected.get('url_pattern', 'reservation')
    if expected_url_pattern.lower() in url.lower():
        score += 0.5
    sort_text = result.get('sort_text', '').lower()
    expected_sort = expected.get('sort_text', '').lower()
    if expected_sort and expected_sort in sort_text:
        score += 0.5
    return min(score, 1.0)

def check_docx_contains__72fc0161fa321ae4bf01e0be489f2375(result, expected, **options):
    """Check if the docx text contains the expected book title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    text = result.get('text', '').strip().lower()
    expected_title = expected.get('expected_title', '').strip().lower()
    if not expected_title:
        return 0.0
    if expected_title in text:
        return 1.0
    words = expected_title.split()
    if not words:
        return 0.0
    matched = sum((1 for w in words if w in text))
    return matched / len(words) * 0.5

def check_docx_table_bold__b70d187abe0811a915df187857b97f14(result, expected, **options):
    """Check if all words in the table are bold. Supports partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    total = result.get('total_words', 0)
    bold = result.get('bold_words', 0)
    if total == 0:
        return 0.0
    if result.get('all_bold', False):
        return 1.0
    return round(bold / total, 2)

def check_docx_contains__896d76d765c6621bc8feff2ba14e6be9(result, expected, **options):
    """Check if the docx text contains the expected book title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    text = result.get('text', '').strip().lower()
    expected_title = expected.get('expected_title', '').strip().lower()
    if not expected_title:
        return 0.0
    if expected_title in text:
        return 1.0
    words = expected_title.split()
    if not words:
        return 0.0
    matched = sum((1 for w in words if w in text))
    return matched / len(words) * 0.5

def check_docx_para_bold__7a1ffa51886994c5f3667d404d37d94d(result, expected, **options):
    """Check if paragraph has bold formatting applied."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    min_ratio = expected.get('min_bold_ratio', 0.9)
    bold_ratio = result.get('bold_ratio', 0.0)
    if bold_ratio >= min_ratio:
        return 1.0
    elif bold_ratio > 0.0:
        return bold_ratio / min_ratio
    return 0.0

def check_docx_text_and_image__2aeb38ddfad7ab2f7bd9e6f63f6cc1bf(result, expected, **options):
    """Check if docx contains expected text and at least one image.

    Partial credit:
      - 0.5 for expected text present (case-insensitive)
      - 0.5 for at least min_images images embedded

    Args:
        result: dict with 'text' (str) and 'image_count' (int) from getter
        expected: dict with 'expected_text' (str) and 'min_images' (int)

    Returns:
        float: 0.0 to 1.0
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    actual_text = result.get('text', '')
    if expected_text and expected_text.lower() in actual_text.lower():
        score += 0.5
    min_images = expected.get('min_images', 1)
    if result.get('image_count', 0) >= min_images:
        score += 0.5
    return min(score, 1.0)

def check_docx_default_font__536a5616674c51027b6cacf28aa7ecd9(result, expected, **options):
    """Check if the default font and run fonts match the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', '')
    score = 0.0
    default_font = result.get('default_font', '')
    if default_font and expected_font.lower() == default_font.lower():
        score += 0.5
    run_fonts = result.get('run_fonts', [])
    if run_fonts:
        matching = sum((1 for f in run_fonts if f.lower() == expected_font.lower()))
        if len(run_fonts) > 0:
            ratio = matching / len(run_fonts)
            score += 0.5 * ratio
    elif score >= 0.5:
        score += 0.5
    return min(score, 1.0)

def check_docx_title__b98b6fc5ce8ebe003cf7eaeacc71d919(result, expected, **options):
    """Check if the docx contains the expected title text."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_title = expected.get('expected_title', '')
    full_text = result.get('text', '')
    if not full_text:
        return 0.0
    if expected_title in full_text:
        return 1.0
    from rapidfuzz import fuzz
    score = fuzz.partial_ratio(expected_title, full_text) / 100.0
    if score >= 0.8:
        return score
    return 0.0

def check_footer_page_numbers__806bfdf7c0ee49dcce08f91a2edad9e6(result, expected, **options):
    """Check if the document has page numbers in footers."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    has_page_number = result.get('has_page_number', False)
    expected_value = expected.get('expected', True)
    return 1.0 if has_page_number == expected_value else 0.0

def check_docx_alignment__ae9bdf892a1bf6b8ab5ef93251d291e0(result, expected, **options):
    """Check if the first paragraph alignment matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('alignment', '')
    expected_alignment = expected.get('expected_alignment', '')
    return 1.0 if actual == expected_alignment else 0.0

def check_title_size_content_bold_bg__80858354e7972b18d1e56bf385980263(result, expected, **options):
    """Check title size, content bold, and background color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_size = expected.get('title_size_pt', 48)
    actual_size = result.get('title_size_pt')
    if actual_size is not None and abs(actual_size - expected_size) < 1.0:
        score += 0.33
    if result.get('content_bold') is True:
        score += 0.33
    expected_bg = expected.get('bg_color', 'FFFF00')
    actual_bg = result.get('bg_color', '')
    if actual_bg:
        actual_upper = actual_bg.upper()
        if expected_bg.upper() in actual_upper or actual_upper in ('FFFF00', 'FFD700'):
            score += 0.34
    return min(score, 1.0)

def check_bg_fontsize__c22d8b15c916d1ef420eb42daeb27bf0(result, expected, **options):
    """Check background color and font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    tolerance = expected.get('color_tolerance', 50)
    expected_bg = expected.get('expected_bg_color', '')
    actual_bg = result.get('bg_color', '')
    color_dist = 999
    if actual_bg and expected_bg:
        try:
            (r1, g1, b1) = (int(actual_bg[0:2], 16), int(actual_bg[2:4], 16), int(actual_bg[4:6], 16))
            (r2, g2, b2) = (int(expected_bg[0:2], 16), int(expected_bg[2:4], 16), int(expected_bg[4:6], 16))
            color_dist = ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
        except Exception:
            color_dist = 999
    if color_dist <= tolerance:
        score += 0.5
    expected_size = expected.get('expected_font_size_pt')
    actual_size = result.get('font_size_pt')
    if expected_size is not None and actual_size is not None:
        if abs(actual_size - expected_size) <= 1.0:
            score += 0.5
    return min(score, 1.0)

def check_heading2_fonts__3ee05560ea0a8eae187c5c1b3f30c5d3(result, expected, **options):
    """Check if all Heading 2 paragraphs use the expected font. Partial credit per heading."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    headings = result.get('headings', [])
    if not headings:
        return 0.0
    expected_font = expected.get('expected_font', 'Arial')
    correct = 0
    for h in headings:
        fonts = h.get('fonts', [])
        if fonts and all((f == expected_font for f in fonts)):
            correct += 1
    return correct / len(headings)

def check_text_underline__1563ed9d69890bfb1ca48bd048aeee4a(result, expected, **options):
    """Check if text runs are all underlined."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_underline = expected.get('expected_underline', True)
    actual_underline = result.get('all_underline', False)
    return 1.0 if actual_underline == expected_underline else 0.0

def check_font_color_match__d0c681af3bd2af47798202cb5ebe1988(result, expected, **options):
    """Check if the font color matches the expected RGB value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_color = result.get('color_rgb', '').upper()
    expected_color = expected.get('expected_color', '').upper()
    if not actual_color or not expected_color:
        return 0.0
    if actual_color == expected_color:
        return 1.0
    try:
        (ar, ag, ab) = (int(actual_color[0:2], 16), int(actual_color[2:4], 16), int(actual_color[4:6], 16))
        (er, eg, eb) = (int(expected_color[0:2], 16), int(expected_color[2:4], 16), int(expected_color[4:6], 16))
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        if distance < 30:
            return 0.5
    except (ValueError, IndexError):
        pass
    return 0.0

def check_docx_para_alignment__099018697289346087fe050b42a1dd16(result, expected, **options):
    """Check alignment of each paragraph. Partial credit per paragraph."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('alignments', [])
    expected_aligns = expected.get('expected_alignments', [])
    if not actual or not expected_aligns:
        return 0.0
    score = 0.0
    n = min(len(actual), len(expected_aligns))
    weight = 1.0 / n if n > 0 else 0.0
    for i in range(n):
        if actual[i] == expected_aligns[i]:
            score += weight
    return min(score, 1.0)

def check_docx_model_names_table__257a7a045a00e35dcfe1e7d02a439bcc(result, expected, **options):
    """Check if docx has a table or text containing both model names from the xlsx.
    Partial credit: 0.5 per model name found in tables or Main Results section text.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_names = expected.get('model_names', ['Gpt-4', 'rl+il'])
    tables = result.get('tables', [])
    section_text = result.get('main_results_text', '')
    all_text = section_text.lower()
    for table in tables:
        for row in table:
            all_text += ' ' + ' '.join(row).lower()
    score = 0.0
    per_name = 1.0 / len(expected_names) if expected_names else 0.0
    for name in expected_names:
        if name.lower() in all_text:
            score += per_name
    return min(score, 1.0)

def check_docx_text_formatting__92a351c44e937610849c0af1deb80b89(result, expected, **options):
    """Check text formatting per paragraph. Partial credit per paragraph."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('paragraphs', [])
    expected_fmt = expected.get('expected_formatting', [])
    if not actual or not expected_fmt:
        return 0.0
    score = 0.0
    n = min(len(actual), len(expected_fmt))
    weight = 1.0 / n if n > 0 else 0.0
    for i in range(n):
        match = True
        for key in expected_fmt[i]:
            if actual[i].get(key) != expected_fmt[i][key]:
                match = False
                break
        if match:
            score += weight
    return min(score, 1.0)

def check_column_header_and_values__23bfc830003e8f489b53a39808f56b59(result, expected, **options):
    """Check column header and sampled values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    expected_header = expected.get('expected_header', '')
    actual_header = result.get('header', '')
    if actual_header and expected_header:
        if str(actual_header).strip().lower() == str(expected_header).strip().lower():
            score += 0.3
    expected_values = expected.get('expected_values', {})
    actual_values = result.get('values', {})
    tolerance = expected.get('tolerance', 0.5)
    if expected_values:
        per_value = 0.7 / len(expected_values)
        for (row_key, exp_val) in expected_values.items():
            act_val = actual_values.get(row_key)
            if act_val is not None and exp_val is not None:
                try:
                    if abs(float(act_val) - float(exp_val)) <= tolerance:
                        score += per_value
                except (TypeError, ValueError):
                    pass
    return min(score, 1.0)

def check_doc_file_list__d6bd2cbacad8b9b51e1e935f00211e4d(result, expected, **options):
    """Check that filelist.txt contains listing of .doc files."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_doc_count', 12)
    if result.get('filelist_exists'):
        score += 0.4
    content = result.get('filelist_content', '')
    if content and result.get('filelist_exists'):
        doc_names_found = 0
        for line in content.split('\n'):
            line = line.strip()
            if line.endswith('.doc') and (not line.endswith('.docx')):
                doc_names_found += 1
        if doc_names_found >= expected_count:
            score += 0.6
        elif doc_names_found > 0:
            score += 0.6 * (doc_names_found / expected_count)
    return min(score, 1.0)

def check_docx_gpt4_specific_scores__9467a3bac5aae80d6f5894cff6144c37(result, expected, **options):
    """Check if docx table contains GPT-4's os and chrome scores.
    Partial credit: 0.34 for Gpt-4 name, 0.33 for os score, 0.33 for chrome score.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    tables = result.get('tables', [])
    if not tables:
        return 0.0
    expected_name = expected.get('model_name', 'Gpt-4').lower()
    expected_os = float(expected.get('os_score', '0.3333'))
    expected_chrome = float(expected.get('chrome_score', '0.3636'))
    score = 0.0
    found_name = False
    found_os = False
    found_chrome = False
    for table in tables:
        for row in table:
            row_lower = ' '.join([cell.lower() for cell in row])
            if expected_name in row_lower:
                found_name = True
            for cell in row:
                cell_clean = cell.strip()
                try:
                    val = float(cell_clean)
                    if abs(val - expected_os) < 0.01:
                        found_os = True
                    if abs(val - expected_chrome) < 0.01:
                        found_chrome = True
                except (ValueError, TypeError):
                    pass
    if found_name:
        score += 0.34
    if found_os:
        score += 0.33
    if found_chrome:
        score += 0.33
    return min(score, 1.0)

def check_font_color__3892513c1cab5055660dfeaa188121fb(result, expected, **options):
    """Check if all text runs have the expected font color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    colors = result.get('colors', [])
    expected_color = expected.get('expected_color', '').upper()
    if not colors:
        return 0.0
    matching = sum((1 for c in colors if c and c.upper() == expected_color))
    return matching / len(colors)

def check_page_numbers__9de26a1b58798be1fa2f4b6b3d9e9ec9(result, expected, **options):
    """Check if page numbers are present in the document footer.

    Scoring: 1.0 if page numbers found, 0.0 otherwise.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    has_page_numbers = result.get('has_page_numbers', False)
    expected_value = expected.get('expected_has_page_numbers', True)
    if has_page_numbers == expected_value:
        return 1.0
    return 0.0

def check_pdf_page_count__ddacf8dd85e44a11a4923ff07ce0abff(result, expected, **options):
    """Check if PDF exists and has expected page count."""
    if result.get('error') or not result.get('exists'):
        return 0.0
    score = 0.0
    score += 0.5
    expected_pages = expected.get('expected_page_count', 0)
    actual_pages = result.get('page_count', 0)
    if actual_pages == expected_pages:
        score += 0.5
    return score

def check_all_fonts_changed__f0d9875ec4c08351c2fb33918a5793ea(result, expected, **options):
    """Check if all fonts in the document match the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    fonts = result.get('fonts', [])
    expected_font = expected.get('expected_font', '')
    if not fonts:
        return 0.0
    if len(fonts) == 1 and fonts[0].lower() == expected_font.lower():
        return 1.0
    matching = sum((1 for f in fonts if f.lower() == expected_font.lower()))
    return matching / len(fonts) if fonts else 0.0

def check_header_corrected__e0d3491cdcd2d9a19af59b5d2c376712(result, expected, **options):
    """Check if the header typo in A1 was corrected from 'tile' to 'title'."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('a1', '')
    expected_val = expected.get('expected_value', 'title')
    if actual is None:
        return 0.0
    if str(actual).strip().lower() == str(expected_val).strip().lower():
        return 1.0
    return 0.0

def check_underline__935b1aaa61bb5f50d9e66c69d9b5829b(result, expected, **options):
    """Check if all text runs are underlined."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('all_underlined', False):
        return 1.0
    total = result.get('total_runs', 0)
    underlined = result.get('underlined_runs', 0)
    if total == 0:
        return 0.0
    return underlined / total

def check_document_font__2ae82910da68002280a1510e1ca61e99(result, expected, **options):
    """Check if all text uses the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    score = 0.0
    if result.get('default_font') == expected_font:
        score += 0.3
    total = result.get('total_runs', 0)
    matching = result.get('matching_runs', 0)
    if total > 0:
        ratio = matching / total
        score += 0.7 * ratio
    elif result.get('default_font') == expected_font:
        score += 0.7
    return min(score, 1.0)

def check_flight_params__15e6e06fadd423d4e6a279cca6f8914d(result, expected, **options):
    """Check flight search URL parameters with partial credit.
    Gives 0.5 for correct fromStation, 0.5 for correct toStation.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_from = expected.get('fromStation', '')
    expected_to = expected.get('toStation', '')
    if expected_from and result.get('fromStation', '') == expected_from:
        score += 0.5
    if expected_to and result.get('toStation', '') == expected_to:
        score += 0.5
    return score

def check_title_font_color__dcc119efe7774eb3168a0af998c68532(result, expected, **options):
    """Check if the title font color matches the expected RGB value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_color = result.get('color_rgb', '').upper()
    expected_color = expected.get('expected_color', '').upper()
    if not actual_color or not expected_color:
        return 0.0
    if actual_color == expected_color:
        return 1.0
    try:
        (ar, ag, ab) = (int(actual_color[0:2], 16), int(actual_color[2:4], 16), int(actual_color[4:6], 16))
        (er, eg, eb) = (int(expected_color[0:2], 16), int(expected_color[2:4], 16), int(expected_color[4:6], 16))
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        if distance < 30:
            return 0.5
    except (ValueError, IndexError):
        pass
    return 0.0

def check_has_page_numbers__ad71c47cc02b102c764d5154ce0d73c0(result, expected, **options):
    """Check if the document has page numbers in the footer."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    if result.get('has_page_numbers', False):
        return 1.0
    return 0.0

def check_font_size__e0ab347b041b8b1f40c9ca0e5bc35d39(result, expected, **options):
    """Check if font size matches expected size."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_size = result.get('font_size_pt')
    expected_size = expected.get('expected_size_pt')
    if actual_size is None or expected_size is None:
        return 0.0
    tolerance = expected.get('tolerance', 1.0)
    if abs(float(actual_size) - float(expected_size)) <= tolerance:
        return 1.0
    return 0.0

def check_italic_color__575ab439795c172e41189face7efb98a(result, expected, **options):
    """Check if all italic text has the expected font color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    italic_colors = result.get('italic_colors', [])
    if not italic_colors:
        return 0.0
    expected_color = expected.get('expected_color', 'FF0000')
    matching = sum((1 for c in italic_colors if c == expected_color))
    return matching / len(italic_colors)

def check_title_font__7f1ba412564f3dc4322d685ce4a3274f(result, expected, **options):
    """Check if the title font matches the expected font name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('expected_font', '')
    font_names = result.get('font_names', [])
    if not font_names or not expected_font:
        return 0.0
    matching = sum((1 for f in font_names if f.lower() == expected_font.lower()))
    return matching / len(font_names) if font_names else 0.0

def check_hint_line_spacing__7fbc318f6c5b72e9f18e5cf28491b9c0(result, expected, **options):
    """Check if all five hint paragraphs have the expected line spacing."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    spacings = result.get('spacings', [])
    expected_spacing = expected.get('expected_spacing', 1.5)
    hint_count = result.get('hint_count', 0)
    if hint_count == 0:
        return 0.0
    correct = 0
    for s in spacings:
        if s is not None and abs(s - expected_spacing) < 0.01:
            correct += 1
    return correct / max(len(spacings), 1)

def check_writer_title_italic__088f05c0fd2c52a92bc3dc8f8cc6719c(result, expected, **options):
    """Check if the title is italic."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_italic = expected.get('expected_italic', True)
    actual_italic = result.get('is_italic', False)
    if actual_italic == expected_italic:
        return 1.0
    return 0.0

def check_title_font__1b8a8292a0fd2eb13bf16e945356170a(result, expected, **options):
    """Check if the title font matches the expected font name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_font = expected.get('expected_font', 'Arial')
    fonts = result.get('fonts', [])
    if not fonts:
        return 0.0
    matching = sum((1 for f in fonts if f.lower() == expected_font.lower()))
    return matching / len(fonts)

def check_startup_and_font__aa9af3d49c1447b85e91f3d729dc4ea1_qw35sft2_1947001b(result, expected, **options):
    """
    Partial-credit metric for two sub-goals:
      0.5 - funbrain.com startup is removed
      0.5 - default font size matches expected (20 = Large)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    restore = result.get('restore_on_startup', -1)
    startup_urls = result.get('startup_urls') or []
    has_funbrain = any(('funbrain' in str(url).lower() for url in startup_urls))
    if not has_funbrain or restore != 4:
        score += 0.5
    expected_size = expected.get('default_font_size', 20)
    actual_size = result.get('default_font_size', 16)
    if actual_size == expected_size:
        score += 0.5
    return score

def check_diamond_campground_page__c0084be8e992e4c5b6a50c34c3f5f437_qw35sft2_c0d65494(result, expected, **options):
    """Check if the active tab shows a Diamond campground page on recreation.gov.

    result: dict from active_tab_info with keys: url, title, content
    expected: dict (unwrapped from rules) with keys:
        url_fragment (str): substring expected in URL
        keyword (str): keyword expected in title or content
    """
    if not isinstance(result, dict):
        return 0.0
    url = result.get('url', '') or ''
    title = result.get('title', '') or ''
    content = result.get('content', '') or ''
    url_fragment = expected.get('url_fragment', 'recreation.gov/camping/campgrounds/')
    keyword = expected.get('keyword', 'Diamond')
    url_ok = url_fragment in url
    keyword_ok = keyword in title or keyword in content
    if url_ok and keyword_ok:
        return 1.0
    elif url_ok:
        return 0.5
    return 0.0

def check_chrome_dnt_and_font__f0db2c12ec96d3956f7c7b7b8c915fc1_qw35sft2_7b230a8b(result, expected, **options):
    """Check DNT is enabled (0.5) and font size matches expected value (0.5)."""
    if result is None or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    if result.get('dnt') == expected.get('dnt', True):
        score += 0.5
    expected_font = expected.get('font_size')
    actual_font = result.get('font_size')
    if expected_font is not None and actual_font is not None:
        try:
            if int(actual_font) == int(expected_font):
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score

def check_font_and_dnt__baedb9f4aec1e7501564ae9f37c335ac_qw35sft2_a1f236ee(result, expected, **options):
    """Check Chrome font size (0.5) and Do Not Track enabled (0.5) with partial credit."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    expected_font_size = expected.get('expected_font_size', 24)
    expected_dnt = expected.get('expected_do_not_track', True)
    actual_font = result.get('font_size')
    if actual_font is not None and int(actual_font) == int(expected_font_size):
        score += 0.5
    actual_dnt = result.get('do_not_track')
    if actual_dnt == expected_dnt:
        score += 0.5
    return score

def check_header_bold__988aa3731e056dbe2f61a28b2135e9ae_qw35sft2_54bf3421(result, expected, **options):
    """Check that both header cells A1 and B1 are bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    a1_bold = result.get('A1_bold', False)
    b1_bold = result.get('B1_bold', False)
    score = 0.0
    if a1_bold:
        score += 0.5
    if b1_bold:
        score += 0.5
    return score

def check_bold_header_and_pdf__86ea3ac794566dae35e2163b9c7d3c94_qw35sft2_0702a49e(result, expected, **options):
    """Partial credit: 1/3 for all header cells bold (A11+B11+C11), 1/3 for fit-to-page enabled, 1/3 for PDF exported."""
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('header_bold', False):
        score += 0.34
    if result.get('fit_to_page', False):
        score += 0.33
    if result.get('pdf_exists', False):
        score += 0.33
    return min(score, 1.0)

def check_bold_and_csv__8f6a3ccc5fb938975de38e8fbcc8f580_qw35sft2_69080047(result, expected, **options):
    """Partial credit: 0.5 for all column A cells bold in xlsx, 0.5 for CSV exported."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('csv_exists') and result.get('csv_row_count', 0) >= 4:
        score += 0.5
    if result.get('all_col_a_bold'):
        score += 0.5
    return score

def check_employee_split_bold_header__fc565b338fbb3db18703aef2eaa623cb_qw35sft2_5a37c450(result, expected, **options):
    """Check that data was split correctly AND header cells B1:D1 are bold.

    Partial credit:
    - 0.5: B2/C2/D2 have correct first employee split values
    - 0.5: Header cells B1, C1, D1 are all bold
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    exp_b2 = expected.get('b2', 'Benedict')
    exp_c2 = expected.get('c2', 'Cucumberbach')
    exp_d2 = expected.get('d2', 'Manager')
    b2_ok = str(result.get('b2') or '').strip() == exp_b2
    c2_ok = str(result.get('c2') or '').strip() == exp_c2
    d2_ok = str(result.get('d2') or '').strip() == exp_d2
    if b2_ok and c2_ok and d2_ok:
        score += 0.5
    if result.get('b1_bold') and result.get('c1_bold') and result.get('d1_bold'):
        score += 0.5
    return min(score, 1.0)

def check_bold_header__771ce45281d75f3370a131f92752ad7a_qw35sft2_36c09f03(result, expected, **options):
    """Check C2 (0.25) and C3 (0.25) cleaned titles correct, and both A1+C1 headers bold (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    c2_actual = result.get('c2_value')
    c2_expected = expected.get('c2_expected')
    if c2_actual is not None and c2_expected is not None:
        if str(c2_actual).strip() == str(c2_expected).strip():
            score += 0.25
    c3_actual = result.get('c3_value')
    c3_expected = expected.get('c3_expected')
    if c3_actual is not None and c3_expected is not None:
        if str(c3_actual).strip() == str(c3_expected).strip():
            score += 0.25
    c1_bold = result.get('c1_bold')
    a1_bold = result.get('a1_bold')
    c1_bold_expected = expected.get('c1_bold_expected', True)
    a1_bold_expected = expected.get('a1_bold_expected', True)
    if c1_bold == c1_bold_expected and a1_bold == a1_bold_expected:
        score += 0.5
    return score

def check_header_bold__636602487d408f3cbef02e77b45f21ff_qw35sft2_0ad6b944(result, expected, **options):
    """Check that zoom was reduced from 200% AND all header cells C4:H4 are bold.

    Scoring breakdown:
      - 0.4 * zoom_ok   : zoom_scale < 200 (agent zoomed out from the original 200%)
      - 0.6 * bold_score: fraction of C4:H4 cells (6 total) that have explicit bold

    zoom_scale == 0 means LibreOffice default (100%), which satisfies zoom < 200.
    font.bold == None (no explicit override) is treated as not bold (is True check).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    cells = ['C4', 'D4', 'E4', 'F4', 'G4', 'H4']
    bold_count = sum((1 for c in cells if result.get(c) is True))
    bold_score = bold_count / len(cells)
    zoom_scale = result.get('zoom_scale', 200)
    zoom_ok = 1.0 if zoom_scale < 200 else 0.0
    return 0.4 * zoom_ok + 0.6 * bold_score

def check_title_underline__42b6fd5c06cd718a123cb3a7892d354e_qw35sft2_3c48d0f0(result, expected, **options):
    """Check if all runs in the slide title are underlined."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('all_underlined') is True else 0.0

def check_title_bold_and_size__d1408966fdb4fd77bea2fd3e21b1ee11_qw35sft2_04364771(result, expected, **options):
    """Check that the title text on slide 1 is bold and has the expected font size (partial credit)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('bold') == expected.get('expected_bold', True):
        score += 0.5
    expected_size = expected.get('expected_size_pt')
    actual_size = result.get('size_pt')
    if expected_size and actual_size and (abs(float(actual_size) - float(expected_size)) < 0.5):
        score += 0.5
    return score

def check_strikethrough_and_italic__be4183a9c6982ba683532f14bc50d4bc_qw35sft2_4f1f568a(result, expected, **options):
    """Check strikethrough on two Finance Meetings items + italic on third. Partial credit."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    state = result
    strike_targets = expected.get('strike_targets', [])
    if strike_targets:
        per_item = 0.67 / len(strike_targets)
        for text in strike_targets:
            item = state.get(text, {})
            if isinstance(item, dict) and item.get('strike') is True:
                score += per_item
    italic_target = expected.get('italic_target', '')
    if italic_target:
        item = state.get(italic_target, {})
        if isinstance(item, dict) and item.get('italic') is True:
            score += 0.33
    return min(round(score, 4), 1.0)

def check_title_bold__78dfac329c5cd8b1c5efec5b79600fdd_qw35sft2_99d8b746(result, expected, **options):
    """Check that the title on slide 2 is bold. Returns 1.0 if bold, 0.0 otherwise."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    return 1.0 if result.get('title_bold', False) is True else 0.0

def check_title_font_size__f083039072ffd98eec2d101bb28501a7_qw35sft2_4d31c951(result, expected, **options):
    """Check that the title text on slide 1 has the expected font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_size_pt')
    actual_size = result.get('size_pt')
    if expected_size is None or actual_size is None:
        return 0.0
    return 1.0 if abs(float(actual_size) - float(expected_size)) < 0.5 else 0.0

def check_writer_footer_body_spacing__e86d756354df667eec5758969f058362_qw35sft2_7eb35045(result, expected, **options):
    """Partial credit: 0.5 for footer PAGE field, 0.5 for 1.5 line spacing on body paragraph."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_page_field'):
        score += 0.5
    if result.get('body_spacing_1_5'):
        score += 0.5
    return score

def check_title_run_font__ac8fd7a53908049060d14f1a620d5f79_qw35sft2_f9a68cbe(result, expected, **options):
    """Check whether the document title's first run uses the expected font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_font = expected.get('font_name', '')
    actual_font = result.get('font_name') or ''
    if not expected_font:
        return 0.0
    return 1.0 if actual_font.strip().lower() == expected_font.strip().lower() else 0.0

def check_para0_underline__0ee57f85e0ff64364a5561938bc94f89_qw35sft2_e008e549(result, expected, **options):
    """Check if all runs in the first paragraph are underlined."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_underline = expected.get('expected_underline', True)
    actual_underline = result.get('underline', False)
    return 1.0 if actual_underline == expected_underline else 0.0

def check_docx_last_line_value__9ae69d0ef002958997c53c817727d757_qw35sft2_933992b3(result, expected, **options):
    """Check if the last line of the docx matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    last_line = str(result.get('last_line', '')).strip()
    expected_value = str(expected.get('expected_value', '')).strip()
    if not expected_value:
        return 0.0
    if last_line == expected_value:
        return 1.0
    if expected_value in last_line:
        return 0.8
    return 0.0

def check_docx_subscript_and_title_italic__5faa01d8a062ee366a1e56b39b81f8d9_qw35sft2_9e3c6cfb(result, expected, **options):
    """
    Check:
    - 0.5 pts: '2' in title has subscript formatting
    - 0.5 pts: 'Soak up the Science' run in title has italic formatting
    expected is the rules dict (already unwrapped by get_rule()).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('subscript_2_in_title') is True:
        score += 0.5
    if result.get('title_soak_italic') is True:
        score += 0.5
    return score

def check_table_7x5_with_headers__a4c465e58afc03255b4d0ddc46d5d525_qw35sft2_2ca96022(result, expected, **options):
    """Check 7x5 table inserted with expected column headers in first row.
    Partial credit: 0.34 for table count, 0.33 for dimensions, 0.33 for all headers matching.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count') == expected.get('table_count', 2):
        score += 0.34
    if result.get('last_rows') == expected.get('last_rows', 5) and result.get('last_cols') == expected.get('last_cols', 7):
        score += 0.33
    expected_headers = expected.get('first_row', [])
    actual_headers = result.get('first_row', [])
    if expected_headers and len(actual_headers) == len(expected_headers):
        matches = sum((1 for exp, act in zip(expected_headers, actual_headers) if exp.lower() in act.lower()))
        if matches == len(expected_headers):
            score += 0.33
    return min(score, 1.0)

def check_writer_basic_fonts__cb1a08b69b2e8e76cab6534a166b3ea6_qw35sft2_0eb16eec(result, expected, **options):
    """
    Partial credit:
    - 0.5 if default_font == expected_default_font
    - 0.5 if heading_font == expected_heading_font
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_default = expected.get('expected_default_font', 'Times New Roman')
    expected_heading = expected.get('expected_heading_font', 'Liberation Serif')
    if result.get('default_font') == expected_default:
        score += 0.5
    if result.get('heading_font') == expected_heading:
        score += 0.5
    return min(score, 1.0)

def check_docx_lower_and_page_nums__4155f473baa624a7c08f0169367bee57_qw35sft2_31960405(result, expected, **options):
    """Partial credit: 0.5 for all-lowercase text, 0.5 for page numbers in footer."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('all_lower', False):
        score += 0.5
    if result.get('has_page_numbers', False):
        score += 0.5
    return score

def check_font_and_align__2b822025bba1f4240036b27d0339b871_qw35sft2_f2ba7288(result, expected, **options):
    """Check font name (0.5) and center-alignment (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('font_name') == expected_font and result.get('all_same_font', False):
        score += 0.5
    if result.get('all_center_aligned', False):
        score += 0.5
    return score

def check_docx_table_italic_col__838d41f5c432af14f1eb478fa956e99e_qw35sft2_7a2cc52c(result, expected, **options):
    """Check that the table exists with correct rows and first-column data cells are italic.
    Partial credit: 0.5 for table with rows, 0.5 for italic first column.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('row_count', 5)
    if result.get('table_count', 0) >= 1 and result.get('row_count', 0) == expected_rows:
        score += 0.5
    if result.get('data_rows_italic', False) == expected.get('data_rows_italic', True):
        score += 0.5
    return score

def check_writer_heading_body_align__d846f8a9d2cb6584815a6f17cfb8e80a_qw35sft2_481a2564(result, expected, **options):
    """Check heading (0.5) and body paragraph (0.5) are both center-aligned."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('heading_centered') is True:
        score += 0.5
    if result.get('body_centered') is True:
        score += 0.5
    return score

def check_docx_first_para_strike__efe155f44671b1c09a26d5fefbd2fc44_qw35sft2_04b52bc6(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_strike = expected.get('has_strikethrough', True)
    actual_strike = result.get('has_strikethrough', False) if isinstance(result, dict) else False
    return 1.0 if actual_strike == expected_strike else 0.0

def check_all_italic_size_14__b8546f99d1a88f06fb0b6ffbfa13cf55_qw35sft2_ff244508(result, expected, **options):
    """Check that all italic runs in the document are 14pt.
    Returns 1.0 if all italic runs are 14pt, 0.0 otherwise.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('all_italic_size_14', False):
        return 0.0
    if result.get('italic_count', 0) == 0:
        return 0.0
    return 1.0

def check_odt_highlight_italic__e4e69693f939d9cc32e39d11fb21f92a_qw35sft2_86960b9e(result, expected, **options):
    """
    Score = 0.5 (no highlights) + 0.5 (title italic).
    expected keys: no_highlights (bool), title_italic (bool)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_highlights', True):
        score += 0.5
    if result.get('title_italic', False):
        score += 0.5
    return min(score, 1.0)

def check_three_para_spacing__e4c36138929fcbc781e762a8996148c1_qw35sft2_f2c7fe41(result, expected, **options):
    """Check line spacing of first three content paragraphs.
    Partial credit: 1/3 per paragraph matching expected spacing.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    for key in ('para0_spacing', 'para2_spacing', 'para4_spacing'):
        if result.get(key) == expected.get(key):
            score += 1.0 / 3.0
    return round(min(score, 1.0), 4)

def check_writer_extended__d40d48ebd32c5cdd693a8c4ab4565d6c_qw35sft2_577878c3(result, expected, **options):
    """Partial credit: Steinberg added (0.33) + placeholder removed (0.33) + title changed (0.34)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_steinberg', False):
        score += 0.33
    if not result.get('has_add_here_marker', True):
        score += 0.33
    expected_title = expected.get('main_title', '')
    actual_title = result.get('main_title', '')
    if expected_title and actual_title.strip() == expected_title.strip():
        score += 0.34
    return min(score, 1.0)

def check_docx_spacing_arial__d472f99c5ab0bbd719c99178119f321f_qw35sft2_1b4c6542(result, expected, **options):
    """Check line spacings (single/double/1.5) and Arial font throughout. 0.25 pts each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('intro_spacing') == 1.0:
        score += 0.25
    if result.get('body_spacing') == 2.0:
        score += 0.25
    if result.get('conclusion_spacing') == 1.5:
        score += 0.25
    if result.get('all_arial'):
        score += 0.25
    return score

def check_writer_break_and_notes__4e9ea3f60c13559abc5abe689f4b9ea2_qw35sft2_d173d01c(result, expected, **options):
    """
    Partial credit: 0.5 for page break inserted, 0.5 for 'Notes' text on blank page.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    min_breaks = expected.get('min_page_breaks', 5)
    if result.get('explicit_page_breaks', 0) >= min_breaks:
        score += 0.5
    if result.get('notes_found', False):
        score += 0.5
    return score

def check_writer_default_and_list_fonts__12134317917b6b593b5731418bb79a41_qw35sft2_f9b4e005(result, expected, **options):
    """
    Partial credit:
    - 0.5 if default_font == expected_default_font
    - 0.5 if list_font == expected_list_font
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_default = expected.get('expected_default_font', 'Times New Roman')
    expected_list = expected.get('expected_list_font', 'Liberation Sans')
    if result.get('default_font') == expected_default:
        score += 0.5
    if result.get('list_font') == expected_list:
        score += 0.5
    return min(score, 1.0)

def check_docx_first_line_contains__d29d0e1e9fda0acc872853f3a63d2906_qw35sft2_6abfbae7(result, expected, **options):
    """Check if the first line of the docx contains the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    first_line = str(result.get('first_line', '')).strip().lower()
    expected_text = str(expected.get('expected_text', '')).strip().lower()
    if not expected_text:
        return 0.0
    if expected_text in first_line:
        return 1.0
    return 0.0

def check_writer_footer_title_align__54873b40ed0174b0628f230c8ad47868_qw35sft2_7c373a94(result, expected, **options):
    """Partial credit: 0.5 for footer PAGE field, 0.5 for centered title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_page_field'):
        score += 0.5
    if result.get('title_centered'):
        score += 0.5
    return score

def check_docx_lower_and_bold_title__5844d45da47ec52046a23219bd5db770_qw35sft2_3c6668d1(result, expected, **options):
    """Partial credit: 0.5 for all-lowercase text, 0.5 for bold first paragraph."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('all_lower', False):
        score += 0.5
    if result.get('first_para_bold', False):
        score += 0.5
    return score

def check_docx_subscript_and_title_center__49f57bac6148724412fcfb3eb162539d_qw35sft2_fa414f8f(result, expected, **options):
    """
    Check:
    - 0.5 pts: '2' in title has subscript formatting
    - 0.5 pts: title paragraph is center-aligned
    expected is the rules dict (already unwrapped by get_rule()).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('subscript_2_in_title') is True:
        score += 0.5
    if result.get('title_centered') is True:
        score += 0.5
    return score

def check_titlecase_and_italic__7a4ad09a1377c8f0eb798dbce87a7dce_qw35sft2_84c30304(result, expected, **options):
    """Check title case (0.5) and title paragraph italic formatting (0.5).
    expected is already the rules dict (unwrapped by get_rule()).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_case_applied') == expected.get('title_case_applied', True):
        score += 0.5
    if result.get('title_italic') == expected.get('title_italic', True):
        score += 0.5
    return score

def check_pdf_and_odt_exists__c65906d6f690144adc013e7fcea2301e_qw35sft2_02517358(result, expected, **options):
    """Partial credit: 0.5 for PDF file existing, 0.5 for ODT file existing."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('pdf_exists', False):
        score += 0.5
    if result.get('odt_exists', False):
        score += 0.5
    return score

def check_docx_table_structure__638d8a63e2ebfdf518128094d93e23f6_qw35sft2_f58f4a90(result, expected, **options):
    """Check that the document contains a table with the expected row count and first cell text.
    Partial credit: 0.5 for table existence, 0.5 for correct row count.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count', 0) >= 1:
        score += 0.5
    expected_rows = expected.get('row_count', 5)
    if result.get('row_count', 0) == expected_rows:
        score += 0.5
    return score

def check_docx_has_image__ed4f0e4477e03c56fe74b28d9d6a3444_qw35sft2_5c2f253c(result, expected, **options):
    """Check that the document has at least min_count inline images."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    min_count = expected.get('min_count', 2)
    actual_count = result.get('image_count', 0)
    return 1.0 if actual_count >= min_count else 0.0

def check_para0_word_replaced__469a4ad5eae55be7063b8ec35b77b37d_qw35sft2_a80ef772(result, expected, **options):
    """Check if 'talkative' was replaced with 'participative' in the first paragraph."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    has_participative = result.get('has_participative', False)
    has_talkative = result.get('has_talkative', True)
    if has_participative and (not has_talkative):
        return 1.0
    return 0.0

def check_docx_multi_para_strike__d8440cb16db3c65fff5a2b530ef73072_qw35sft2_d2ea5e29(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('first_para_strike') == expected.get('first_para_strike', True):
        score += 0.5
    if result.get('last_para_strike') == expected.get('last_para_strike', True):
        score += 0.5
    return score

def check_italic_size_and_underline__dce073995a5927fc743181d7c02c0659_qw35sft2_c208bbd1(result, expected, **options):
    """Check two objectives with partial credit:
    - 0.5: All italic runs are 14pt.
    - 0.5: All italic runs are underlined.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('italic_count', 0) == 0:
        return 0.0
    score = 0.0
    if result.get('all_italic_size_14', False):
        score += 0.5
    if result.get('all_italic_underlined', False):
        score += 0.5
    return score

def check_font_and_size__a567f4bdee19e51d5ed6a587d8d38f9d_qw35sft2_224f4149(result, expected, **options):
    """Check font name (0.5) and font size (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    expected_size = float(expected.get('expected_size_pt', 14.0))
    if result.get('font_name') == expected_font and result.get('all_same_font', False):
        score += 0.5
    actual_size = result.get('font_size_pt')
    if actual_size is not None and abs(float(actual_size) - expected_size) < 0.5 and result.get('all_same_size', False):
        score += 0.5
    return score

def check_writer_heading_align_size__a3c76448cfe24f131aa7bba8054b5d03_qw35sft2_4fa6f53a(result, expected, **options):
    """Check heading is center-aligned (0.5) and font size matches expected (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_centered') is True:
        score += 0.5
    expected_size = float(expected.get('expected_font_size_pt', 14))
    font_size = result.get('font_size_pt')
    if font_size is not None and abs(float(font_size) - expected_size) < 0.5:
        score += 0.5
    return score

def check_odt_highlight_font__582984f0118e076c9ea08b5f36393d3b_qw35sft2_d4cd9443(result, expected, **options):
    """
    Score = 0.5 (no highlights) + 0.5 (default/document font matches expected_font).
    expected keys: expected_font (str, case-insensitive match)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_highlights', True):
        score += 0.5
    expected_font = expected.get('expected_font', 'Arial').strip().lower()
    actual_font = (result.get('default_font') or '').strip().lower()
    if expected_font and actual_font and (expected_font in actual_font):
        score += 0.5
    return min(score, 1.0)

def check_writer_three_goals__c8e0b7dd7367f091fd322e812e9986d0_qw35sft2_2a25ad24(result, expected, **options):
    """Partial credit: Steinberg added (0.34) + placeholder removed (0.33) + heading renamed (0.33)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_steinberg', False):
        score += 0.34
    if not result.get('has_add_here_marker', True):
        score += 0.33
    expected_heading = expected.get('refs_heading', 'Reference List')
    actual_heading = result.get('refs_heading', '')
    if actual_heading.strip() == expected_heading.strip():
        score += 0.33
    return min(score, 1.0)

def check_docx_spacing_fontsize__122cebb30dd4bbe8d6067021381568cf_qw35sft2_c1d5de85(result, expected, **options):
    """Check line spacings (single/double/1.5) and font size 14pt. 0.25 pts each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('intro_spacing') == 1.0:
        score += 0.25
    if result.get('body_spacing') == 2.0:
        score += 0.25
    if result.get('conclusion_spacing') == 1.5:
        score += 0.25
    if result.get('all_14pt'):
        score += 0.25
    return score

def check_writer_break_and_font__52445c32d012b75894f7e753d9bb73ed_qw35sft2_a130d84f(result, expected, **options):
    """
    Partial credit: 0.5 for page break inserted, 0.5 for title font changed to expected_font.
    Font name comparison is case-insensitive.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    min_breaks = expected.get('min_page_breaks', 5)
    if result.get('explicit_page_breaks', 0) >= min_breaks:
        score += 0.5
    expected_font = expected.get('expected_font', '')
    actual_font = result.get('title_font') or ''
    if expected_font and expected_font.lower() == actual_font.lower():
        score += 0.5
    return score

def check_titlecase_and_underline__d22455594f3cb66658bea4571465f8f1_qw35sft2_41d40f3d(result, expected, **options):
    """Check title case (0.5) and title paragraph underline formatting (0.5).
    expected is already the rules dict (unwrapped by get_rule()).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_case_applied') == expected.get('title_case_applied', True):
        score += 0.5
    if result.get('title_underline') == expected.get('title_underline', True):
        score += 0.5
    return score

def check_writer_footer_pageno__035b6f0ae922d6e2d1f24326ec60c904_qw35sft2_be318513(result, expected, **options):
    """Check that the footer contains a PAGE field. Returns 1.0 if present, 0.0 otherwise."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('has_page_field') else 0.0

def check_docx_text_all_upper__cbc878ff0e7a736898047bffbd6f85f6_qw35sft2_32945c7e(result, expected, **options):
    """Return 1.0 if all alphabetic characters in the document are uppercase."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('all_upper', False) else 0.0

def check_para0_bold__42cf2947548fcbcc72e5e44a50eb60dd_qw35sft2_2fb424d7(result, expected, **options):
    """Check if all runs in the first paragraph are bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('bold', False)
    return 1.0 if actual_bold == expected_bold else 0.0

def check_docx_train_removed__aa11b4f1fde2d6cc216fd8dac61371d5_qw35sft2_742684e6(result, expected, **options):
    """Check that a specific train ID has been removed from the docx and total lines match."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('target_count', -1) == 0:
        score += 0.6
    expected_total = expected.get('expected_total_lines')
    if expected_total is not None and result.get('total_lines') == expected_total:
        score += 0.4
    return min(score, 1.0)

def check_writer_font_and_alignment__70647269807c892c9cd9454bb994336f_qw35sft2_b7538338(result, expected, **options):
    """
    Partial credit:
    - 0.5 if default_font == expected_font
    - 0.5 if first_para_centered == True
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('default_font') == expected_font:
        score += 0.5
    if result.get('first_para_centered'):
        score += 0.5
    return min(score, 1.0)

def check_image_and_fontsize__285ee55215732bb31d8255226d83a755_qw35sft2_8e5c91c4(result, expected, **options):
    """Partial credit: 0.5 for image inserted, 0.5 for Figure 1 caption at expected font size."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    min_count = expected.get('min_image_count', 2)
    if result.get('image_count', 0) >= min_count:
        score += 0.5
    expected_pt = expected.get('caption_font_size_pt', 12)
    if result.get('caption_font_size_pt') == expected_pt:
        score += 0.5
    return score

def check_docx_table_header_align__83874f16b468982748e51b1480830a3a_qw35sft2_1e849fbb(result, expected, **options):
    """Check that the table exists and header row cells are center-aligned.
    Partial credit: 0.5 for table with rows, 0.5 for center-aligned header.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('row_count', 5)
    if result.get('table_count', 0) >= 1 and result.get('row_count', 0) == expected_rows:
        score += 0.5
    if result.get('header_centered', False) == expected.get('header_centered', True):
        score += 0.5
    return score

def check_docx_subscript_and_bold_heading__fa57f784af7171fa0e677f21369b4177_qw35sft2_2b7094fd(result, expected, **options):
    """
    Check:
    - 0.5 pts: '2' in title has subscript formatting
    - 0.5 pts: 'Fact sheet' heading text is bold
    expected is the rules dict (already unwrapped by get_rule()).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('subscript_2_in_title') is True:
        score += 0.5
    if result.get('bold_heading') is True:
        score += 0.5
    return score

def check_all_italic_size_16__e2de398ff1a46d21a89c93f18cb16653_qw35sft2_97c614d6(result, expected, **options):
    """Check that all italic runs in the document are 16pt.
    Returns 1.0 if all italic runs are 16pt, 0.0 otherwise.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('all_italic_size_16', False):
        return 0.0
    if result.get('italic_count', 0) == 0:
        return 0.0
    return 1.0

def check_font_and_bold__3a9eb9ceb78a9dd18ebdb09636bdbe29_qw35sft2_c7731e94(result, expected, **options):
    """Check font name (0.5) and all-bold state (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('font_name') == expected_font and result.get('all_same_font', False):
        score += 0.5
    if result.get('all_bold', False):
        score += 0.5
    return score

def check_writer_heading_align_font__fb136ac47c87dcc7e787c09564b245bc_qw35sft2_5267f139(result, expected, **options):
    """Check heading is center-aligned (0.5) and font name matches expected (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_centered') is True:
        score += 0.5
    expected_font = expected.get('expected_font_name', 'Times New Roman')
    actual_font = result.get('font_name') or ''
    if expected_font.lower() in actual_font.lower():
        score += 0.5
    return score

def check_docx_strike_bold_chain__30608d866c349bd16efe53a1c72a27d7_qw35sft2_dde7a53f(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('last_para_strike') == expected.get('last_para_strike', True):
        score += 0.5
    if result.get('first_para_bold') == expected.get('first_para_bold', True):
        score += 0.5
    return score

def check_odt_highlight_strikethrough__a30ceeb68a8eeb0b8a512ab61a2387e1_qw35sft2_e68b1a1e(result, expected, **options):
    """
    Score = 0.5 (no highlights) + 0.5 (last sentence has strikethrough).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_highlights', True):
        score += 0.5
    if result.get('last_sentence_strikethrough', False):
        score += 0.5
    return min(score, 1.0)

def check_writer_ref_added__dab123b7868d167252ed777095add0ba_qw35sft2_57d735cd(result, expected, **options):
    """Check that Steinberg reference is added (0.5) and <add here> marker is removed (0.5)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_steinberg', False):
        score += 0.5
    if not result.get('has_add_here_marker', True):
        score += 0.5
    return min(score, 1.0)

def check_docx_spacing_bold_intro__7ab11a38694adeb4332550a69fd5abfe_qw35sft2_646f94b4(result, expected, **options):
    """Check line spacings (single/double/1.5) and bold introduction. 0.25 pts each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('intro_spacing') == 1.0:
        score += 0.25
    if result.get('body_spacing') == 2.0:
        score += 0.25
    if result.get('conclusion_spacing') == 1.5:
        score += 0.25
    if result.get('intro_bold'):
        score += 0.25
    return score

def check_writer_break_and_title_italic__0f7fdf2d11eca0ddff6ddca8ba1c7a92_qw35sft2_0d4a6598(result, expected, **options):
    """
    Partial credit: 0.5 for page break inserted, 0.5 for title italic applied.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    min_breaks = expected.get('min_page_breaks', 5)
    if result.get('explicit_page_breaks', 0) >= min_breaks:
        score += 0.5
    if result.get('title_italic', False) is True:
        score += 0.5
    return score

def check_docx_sentence_case__121c5b95c7fa376e725e2a74f1f18a20_qw35sft2_fe4e5dd5(result, expected, **options):
    """Return 1.0 if the document text follows sentence case throughout.

    Sentence case: first alphabetic character of each paragraph/cell block is
    uppercase; all other alphabetic characters are lowercase.
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('is_sentence_case', False) else 0.0

def check_writer_footer_title_italic__6a6cc374d796e945ca1edd575879589a_qw35sft2_1abc48e5(result, expected, **options):
    """Partial credit: 0.5 for footer PAGE field, 0.5 for italic document title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_page_field'):
        score += 0.5
    if result.get('title_italic'):
        score += 0.5
    return score

def check_docx_last_line_value__99126397762866f43ef4e6109070d76d_qw35sft2_f21f6be2(result, expected, **options):
    """Check if the last line of the docx matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    last_line = str(result.get('last_line', '')).strip()
    expected_value = str(expected.get('expected_value', '')).strip()
    if not expected_value:
        return 0.0
    if last_line == expected_value:
        return 1.0
    if expected_value in last_line:
        return 0.8
    return 0.0

def check_docx_last_line_value__e6176eae3a47bbdd8f12a06b8b082ed4_qw35sft2_5a3ab949(result, expected, **options):
    """Check if the last line of the docx matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    last_line = str(result.get('last_line', '')).strip()
    expected_value = str(expected.get('expected_value', '')).strip()
    if not expected_value:
        return 0.0
    if last_line == expected_value:
        return 1.0
    if expected_value in last_line:
        return 0.8
    return 0.0

def check_writer_font_and_pagebreak__e2b27ec72c919c5a4bd4f151e7ff64b6_qw35sft2_99965063(result, expected, **options):
    """
    Partial credit:
    - 0.5 if default_font == expected_font
    - 0.5 if page_break_count >= 1
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('default_font') == expected_font:
        score += 0.5
    if result.get('page_break_count', 0) >= 1:
        score += 0.5
    return min(score, 1.0)

def check_docx_table_bold_header__d0985f126b123d58872633f0350725ab_qw35sft2_a8a693fa(result, expected, **options):
    """Check that a table was created and its first row is bold.
    Partial credit: 0.5 for table with correct rows, 0.5 for bold header row.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('row_count', 5)
    if result.get('table_count', 0) >= 1 and result.get('row_count', 0) == expected_rows:
        score += 0.5
    if result.get('header_bold', False) == expected.get('header_bold', True):
        score += 0.5
    return score

def check_docx_subscript_title_and_body__cb7672ac12071a9e4c398d11dc4c470e_qw35sft2_7e8df99d(result, expected, **options):
    """
    Check:
    - 0.5 pts: '2' in title paragraph has subscript formatting
    - 0.5 pts: '2' in first body paragraph's 'H2O—SOAK UP THE SCIENCE' has subscript
    expected is the rules dict (already unwrapped by get_rule()).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('subscript_in_title') is True:
        score += 0.5
    if result.get('subscript_in_body') is True:
        score += 0.5
    return score

def check_font_and_italic__c3632663016de6f20637b59e83145d05_qw35sft2_44c1680c(result, expected, **options):
    """Check font name (0.5) and all-italic state (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('font_name') == expected_font and result.get('all_same_font', False):
        score += 0.5
    if result.get('all_italic', False):
        score += 0.5
    return score

def check_doc_page_breaks__b7c367c074a6362d7d9c85a08867a367_qw35sft2_f869087d(result, expected, **options):
    """Check whether the document contains a page break specifically before the 'Conclusion' section."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_value = expected.get('page_break_before_conclusion', True)
    actual_value = result.get('page_break_before_conclusion', False)
    return 1.0 if actual_value == expected_value else 0.0

def check_italic_and_title_size__792e97b1e8b831bfb11064f6cde55b00_qw35sft2_dc097e8b(result, expected, **options):
    """Check two objectives with partial credit:
    - 0.5: All italic runs are 14pt.
    - 0.5: Title (first paragraph) runs are 16pt.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('italic_count', 0) > 0 and result.get('all_italic_size_14', False):
        score += 0.5
    if result.get('title_run_count', 0) > 0 and result.get('all_title_size_16', False):
        score += 0.5
    return score

def check_writer_heading_align_italic__edb463affe3b938af85226bb52cb1b03_qw35sft2_e5e99d63(result, expected, **options):
    """Check heading is center-aligned (0.5) and italic (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_centered') is True:
        score += 0.5
    if result.get('is_italic') is True:
        score += 0.5
    return score

def check_para0_italic__d2318c984412daccf01ad6c479224170_qw35sft2_d6dcf621(result, expected, **options):
    """Check if all runs in the first paragraph are italic."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_italic = expected.get('expected_italic', True)
    actual_italic = result.get('italic', False)
    return 1.0 if actual_italic == expected_italic else 0.0

def check_docx_last_para_bold__c977b8be915d8296c568265b630cf188_qw35sft2_909f71e3(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_bold = expected.get('has_bold', True)
    actual_bold = result.get('has_bold', False) if isinstance(result, dict) else False
    return 1.0 if actual_bold == expected_bold else 0.0

def check_docx_spacing_italic_conclusion__bb617d210691b208d4c61100467118c1_qw35sft2_c1ca409b(result, expected, **options):
    """Check line spacings (single/double/1.5) and italic conclusion paragraph. 0.25 pts each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('intro_spacing') == 1.0:
        score += 0.25
    if result.get('body_spacing') == 2.0:
        score += 0.25
    if result.get('conclusion_spacing') == 1.5:
        score += 0.25
    if result.get('conclusion_italic'):
        score += 0.25
    return score

def check_writer_ref_footer__78615d30a3700f795271d9e8e63e0686_qw35sft2_4066d575(result, expected, **options):
    """Partial credit: Steinberg added (0.4) + placeholder removed (0.3) + footer page num (0.3)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_steinberg', False):
        score += 0.4
    if not result.get('has_add_here_marker', True):
        score += 0.3
    if result.get('has_footer_page_num', False):
        score += 0.3
    return min(score, 1.0)

def check_writer_page_break_count__e042d7b442613a05635c554b051c41b1_qw35sft2_d2800da9(result, expected, **options):
    """Return 1.0 if the document has at least min_page_breaks explicit page breaks."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    min_breaks = expected.get('min_page_breaks', 5)
    actual = result.get('explicit_page_breaks', 0)
    return 1.0 if actual >= min_breaks else 0.0

def check_odt_highlight_fontsize__4ce1d37833ab4f6aa409d7454a343799_qw35sft2_1f515809(result, expected, **options):
    """
    Score = 0.5 (no highlights) + 0.5 (title font size >= expected_font_size).
    expected keys: expected_font_size (number, default 16)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if not result.get('has_highlights', True):
        score += 0.5
    expected_size = float(expected.get('expected_font_size', 16))
    actual_size = result.get('title_font_size')
    if actual_size is not None and actual_size >= expected_size:
        score += 0.5
    return min(score, 1.0)

def check_line_spacing__80aefd2088f27d1be89724efb49ed091_qw35sft2_c6d4321c(result, expected, **options):
    """Check that all non-empty paragraphs have the expected line spacing (as a multiple)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required = float(expected.get('line_spacing', 2.0))
    tolerance = float(expected.get('tolerance', 0.15))
    spacings = result.get('spacings', [])
    if not spacings:
        return 0.0
    matching = sum((1 for s in spacings if s.get('line_spacing') is not None and abs(s['line_spacing'] - required) <= tolerance))
    return matching / len(spacings)

def check_writer_footer_and_header__5d6628525b72a5805637db6ac940e950_qw35sft2_b0dc9230(result, expected, **options):
    """Partial credit: 0.5 for footer PAGE field, 0.5 for non-empty header."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_page_field'):
        score += 0.5
    if result.get('has_header_text'):
        score += 0.5
    return score

def check_pdf_and_footer_pagenum__0bb862e1956085267ced241046059b10_qw35sft2_aca248a7(result, expected, **options):
    """Partial credit: 0.5 for PDF file existing, 0.5 for page numbers in footer."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('pdf_exists', False):
        score += 0.5
    if result.get('has_page_numbers', False):
        score += 0.5
    return score

def check_image_and_pagenumbers__6390ff8ee787412726544ea6ba43db0d_qw35sft2_ee2dbcca(result, expected, **options):
    """Partial credit: 0.5 for image inserted, 0.5 for footer page numbers present."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    min_count = expected.get('min_image_count', 2)
    if result.get('image_count', 0) >= min_count:
        score += 0.5
    if result.get('has_page_numbers') == expected.get('has_page_numbers', True):
        score += 0.5
    return score

def check_docx_subscript_and_heading_underline__b75f28775a1b3fe6faeb633d14a05fab_qw35sft2_266881b5(result, expected, **options):
    """
    Check:
    - 0.5 pts: '2' in title has subscript formatting
    - 0.5 pts: 'Fact sheet' heading text has underline formatting
    expected is the rules dict (already unwrapped by get_rule()).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('subscript_2_in_title') is True:
        score += 0.5
    if result.get('heading_underlined') is True:
        score += 0.5
    return score

def check_titlecase_and_bold__2777e85b511814778d7406b21395647f_qw35sft2_af2391c1(result, expected, **options):
    """Check title case (0.5) and title paragraph bold (0.5).
    expected is already the rules dict (unwrapped by get_rule()).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_case_applied') == expected.get('title_case_applied', True):
        score += 0.5
    if result.get('title_bold') == expected.get('title_bold', True):
        score += 0.5
    return score

def check_footer_page_numbers__699ab5651848f548e06466de9777d875_qw35sft2_417f7c36(result, expected, **options):
    """Check whether the document footer contains page number fields."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_value = expected.get('has_page_numbers', True)
    actual_value = result.get('has_page_numbers', False)
    return 1.0 if actual_value == expected_value else 0.0

def check_docx_dedup__f922eeeed3d49013fd1e13103dbfb120_qw35sft2_9b95ea64(result, expected, **options):
    """Check that the docx has been deduplicated: each train ID appears exactly once."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('all_unique', False):
        score += 0.5
    expected_count = expected.get('expected_unique_count', 8)
    if result.get('unique_train_count') == expected_count and result.get('total_lines') == expected_count:
        score += 0.5
    return min(score, 1.0)

def check_writer_font_and_footer__d1ae06a2cb3f08c6662cf614ce58e285_qw35sft2_70843d4e(result, expected, **options):
    """
    Partial credit:
    - 0.5 if default_font == expected_font
    - 0.5 if has_page_numbers == True
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    if result.get('default_font') == expected_font:
        score += 0.5
    if result.get('has_page_numbers'):
        score += 0.5
    return min(score, 1.0)

def check_italic_size_and_bold__57bf430ce2e461a40bf234942a8ba4ad_qw35sft2_5635b7b0(result, expected, **options):
    """Check that all italic runs are 14pt (0.5 credit) and bold (0.5 credit).
    Partial credit: 0.5 per satisfied sub-goal.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('italic_count', 0) == 0:
        return 0.0
    score = 0.0
    if result.get('all_italic_size_14', False):
        score += 0.5
    if result.get('all_italic_bold', False):
        score += 0.5
    return score

def check_font_and_size__5934b0de41c172c0d1662adc3bbc874f_qw35sft2_1bc37137(result, expected, **options):
    """Check font name (0.5) and font size (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_font = expected.get('expected_font', 'Times New Roman')
    expected_size = float(expected.get('expected_size_pt', 12.0))
    if result.get('font_name') == expected_font and result.get('all_same_font', False):
        score += 0.5
    actual_size = result.get('font_size_pt')
    if actual_size is not None and abs(float(actual_size) - expected_size) < 0.5 and result.get('all_same_size', False):
        score += 0.5
    return score

def check_writer_heading_align_body_bold__e1c39faa270c25698b597d87598bba49_qw35sft2_9c96e15c(result, expected, **options):
    """Check heading is center-aligned (0.5) and body paragraph is bold (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('heading_centered') is True:
        score += 0.5
    if result.get('body_bold') is True:
        score += 0.5
    return score

def check_para0_font_size__7211ff5a62d6ff910c67632b31846d1c_qw35sft2_aa8614af(result, expected, **options):
    """Check if the first paragraph's font size matches the expected point size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_pt = expected.get('expected_pt', 14)
    actual_pt = result.get('font_size_pt')
    if actual_pt is None:
        return 0.0
    return 1.0 if abs(actual_pt - expected_pt) < 0.5 else 0.0

def check_docx_table_last_row__e0b81f05c0001516085ceb6f32f7b48c_qw35sft2_bc759d94(result, expected, **options):
    """Check that the table has the expected row count and the last row's first cell matches.
    Partial credit: 0.5 for row count, 0.5 for last row content.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('row_count', 6)
    if result.get('table_count', 0) >= 1 and result.get('row_count', 0) == expected_rows:
        score += 0.5
    expected_last = expected.get('last_row_cell0', '5 letters making 1 sound')
    actual_last = result.get('last_row_cell0', '')
    if expected_last.lower() in actual_last.lower() or actual_last.lower() in expected_last.lower():
        score += 0.5
    return score

def check_docx_second_para_italic__aa396abb2d9562e18e97067a7c8c6fe9_qw35sft2_83fb48d7(result, expected, **options):
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_italic = expected.get('has_italic', True)
    actual_italic = result.get('has_italic', False) if isinstance(result, dict) else False
    return 1.0 if actual_italic == expected_italic else 0.0

def check_docx_spacing_center_intro__5f00d786b29854e31074fa3cfeefea1b_qw35sft2_55dd34c9(result, expected, **options):
    """Check line spacings (single/double/1.5) and centered introduction. 0.25 pts each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('intro_spacing') == 1.0:
        score += 0.25
    if result.get('body_spacing') == 2.0:
        score += 0.25
    if result.get('conclusion_spacing') == 1.5:
        score += 0.25
    if result.get('intro_centered'):
        score += 0.25
    return score

def check_writer_citation14__21b44c209dfc410a39a6375bc0042826_qw35sft2_2c6f58ae(result, expected, **options):
    """Check: Steinberg reference added (0.5) AND [14] citation appears in body text (0.5)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_steinberg', False):
        score += 0.5
    if result.get('has_bracket14_in_body', False) and (not result.get('has_add_here_marker', True)):
        score += 0.5
    return min(score, 1.0)

def check_docx_with_header__2354b786c1a1e08afc94d3d722cbb7a6_qw35sft2_f6f184cb(result, expected, **options):
    """
    Partial credit scoring:
    - 0.4: docx starts with the expected header line
    - 0.6: docx contains all expected notes after the header
    """
    if result.get('error') or not result.get('lines'):
        return 0.0
    score = 0.0
    lines = result['lines']
    expected_header = expected.get('header', 'Slide Notes')
    expected_notes = expected.get('expected_notes', [])
    first_line = result.get('first_line', '')
    if first_line and expected_header.strip().lower() in first_line.strip().lower():
        score += 0.4
    if expected_notes:
        matched = 0
        for note in expected_notes:
            for line in lines[1:]:
                if note.strip().lower() in line.strip().lower():
                    matched += 1
                    break
        score += 0.6 * (matched / len(expected_notes))
    return round(min(score, 1.0), 2)

def check_docx_title__40a84e6a55059151959a2b459d4c099e_qw35sft2_3441cc96(result, expected, **options):
    """Check if the first paragraph of a docx matches the expected title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('first_paragraph', '') if isinstance(result, dict) else str(result)
    expected_title = expected.get('expected_title', '')
    if expected_title and actual.strip() == expected_title.strip():
        return 1.0
    return 0.0

def check_python_docs_large_font__492670c8dd0e3374d756c80fee290bb9_qw35sft2_4ecbb99e(result, expected, **options):
    """Check if Python docs are open AND Chrome font size is set to Large (20).

    result: dict with 'active_url' and 'font_size' keys
    expected: dict with 'expected_url' and 'expected_font_size' keys
    Returns: 0.5 per satisfied condition, max 1.0
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_url = expected.get('expected_url', 'docs.python.org')
    expected_font = expected.get('expected_font_size', 20)
    active_url = result.get('active_url', '')
    font_size = result.get('font_size')
    score = 0.0
    if expected_url in active_url:
        score += 0.5
    if font_size is not None and int(font_size) == int(expected_font):
        score += 0.5
    return score

def check_docx_duration__48db15ca4dac03d9c5bae1d76e883852_qw35sft2_7772db27(result, expected, **options):
    """Check if the duration line in a docx matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    duration_line = result.get('duration_line', '') if isinstance(result, dict) else str(result)
    expected_duration = expected.get('expected_duration', '')
    if not expected_duration:
        return 0.0
    if expected_duration.lower() in duration_line.lower():
        return 1.0
    return 0.0

def check_docx_gemini_paragraphs__abfc7551de3c0f45f173130d54033329_qw35sft2_d5a836b4(result, expected, **options):
    """Check that gemini_results.docx has at least min_paragraphs non-empty paragraphs
    and that required_substrings appear in the full text."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    min_paragraphs = expected.get('min_paragraphs', 6)
    required_substrings = expected.get('required_substrings', [])
    if result.get('paragraph_count', 0) >= min_paragraphs:
        score += 0.5
    full_text = result.get('full_text', '')
    if required_substrings:
        matches = sum((1 for s in required_substrings if s in full_text))
        score += 0.5 * (matches / len(required_substrings))
    elif full_text:
        score += 0.5
    return min(score, 1.0)

def check_pandoc_installed__79ba7b4d76657204de8c0df9dc21c02e_qw35sft2_b9982ce5(result, expected, **options):
    """Return 1.0 if pandoc is installed and reachable via PATH."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error') and (not result.get('installed')):
        return 0.0
    return 1.0 if result.get('installed') else 0.0

def check_vscode_theme_font_wrap__12b042cae58f794626351d3c568391b7_qw35sft2_b1232fc0(result, expected, **options):
    """Check colorTheme (0.34), editor.fontSize (0.33), editor.wordWrap (0.33) with partial credit."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    if result.get('colorTheme') == expected.get('colorTheme'):
        score += 0.34
    font_expected = expected.get('fontSize')
    font_actual = result.get('fontSize')
    try:
        if font_expected is not None and font_actual is not None and (float(font_actual) == float(font_expected)):
            score += 0.33
    except (TypeError, ValueError):
        pass
    if result.get('wordWrap') == expected.get('wordWrap'):
        score += 0.33
    return min(score, 1.0)

def check_ext_and_fontsize__fe14a817663aeade8a20cb0f5baad2b6_qw35sft2_9dd1b11f(result, expected, **options):
    """Check autoDocstring installed (0.5) + editor.fontSize matches (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'njpwerner.autodocstring')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    actual_font = result.get('font_size')
    expected_font = expected.get('font_size')
    if expected_font is not None and actual_font is not None:
        try:
            if int(actual_font) == int(expected_font):
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score

def check_workspace_and_fontsize__4cb0241f3e1c4176b65966dc96627af4_qw35sft2_c55e228e(result, expected, **options):
    """Partial credit: 0.5 for workspace saved, 0.5 for correct fontSize."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_font_size = expected.get('expected_font_size', 16)
    if result.get('workspace_exists', False):
        score += 0.5
    actual_font = result.get('editor_font_size')
    if actual_font is not None and int(actual_font) == int(expected_font_size):
        score += 0.5
    return score

def check_vscode_debug_focus_font__802820e3fd86caf9b92a9751803125a2_qw35sft2_84c57be3(result, expected, **options):
    """Check debug.focusEditorOnBreak is False and editor.fontSize matches. 0.5 credit each."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('focus_editor_on_break') == expected.get('focus_editor_on_break'):
        score += 0.5
    expected_font = expected.get('font_size')
    actual_font = result.get('font_size')
    if expected_font is not None and actual_font == expected_font:
        score += 0.5
    return min(score, 1.0)

def check_ext_and_fontsize__bb23bc65a82d76a52cba97bb0a6bf9bc_qw35sft2_439af9b1(result, expected, **options):
    """Check extension installed (0.5) + editor.fontSize matches (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'undefined_publisher.test')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    actual_font = result.get('font_size')
    expected_font = expected.get('font_size')
    if expected_font is not None and actual_font is not None:
        try:
            if int(actual_font) == int(expected_font):
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score

def check_vscode_exclude_fontsize__33cb4b35f94afe16fa1fe71880064b0b_qw35sft2_0eef60d3(result, expected, **options):
    """Check files.exclude has pycache pattern (0.5) and editor.fontSize matches (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    files_exclude = result.get('files_exclude', {})
    pycache_pattern = expected.get('pycache_pattern', '**/__pycache__')
    if files_exclude.get(pycache_pattern):
        score += 0.5
    expected_font = expected.get('font_size')
    actual_font = result.get('font_size')
    if expected_font is not None and actual_font is not None:
        try:
            if int(actual_font) == int(expected_font):
                score += 0.5
        except (TypeError, ValueError):
            pass
    return score
