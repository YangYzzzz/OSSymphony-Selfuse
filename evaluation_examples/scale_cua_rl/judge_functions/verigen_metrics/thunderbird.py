"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_thunderbird_identity__4a5ace1e770bd482f5a0433cd05f4404', 'check_thunderbird_smtp_full__2c49114272e52893d5a44b081f8b14e0', 'check_tb_folder_created__4a5755b15fd8848dbc70d01e3dbb4e89', 'check_thunderbird_folder_state__168d8b8e9196711d91fc3322429f68b5', 'check_tb_msg_count__c7401fcb5b9f30e5bd0dd82a92e3f2df', 'check_tb_email_moved_to_bills__d9f2b609b5f34a856608f932bab06223', 'check_thunderbird_account_config__20583839ff43876dc06848262d06ce9a', 'check_thunderbird_smtp_config__22e52b73b299835b6ec6765d839ed989', 'check_thunderbird_local_folder__4de43750fd473996ecf007ebf44296c8', 'check_git_email_and_push__0c30b3c3ecef4bb03c9efc8f2c5c6e42_qw35sft2_f4164a4f', 'check_tb_dark_and_reply__43d8038e02a0c2b1b0f063b3804b00e6_qw35sft2_12b5fb73', 'check_thunderbird_sig_name__840f3ea7ded7f1b6695bc630a0a035b7_qw35sft2_00505ae2', 'check_thunderbird_local_folders__1ea15ae7572bc843dca8cde6cbcd3a7c_qw35sft2_eaf2591f', 'check_tb_forward_filter_state__f7fbe032caaa49b6e5ebdeafec739cf8_qw35sft2_66712d03', 'check_thunderbird_smtp_no_incoming__f562ce0912b384f5748be41ba253880a_qw35sft2_38a51261', 'check_tb_acct_folder__2e9b3985bbecaa462dbe43eedf3d22fe_qw35sft2_82ec5e33', 'check_thunderbird_draft_cc_attachment__ae64d5efe053bb1ed61ef261c2f7eddd_qw35sft2_cb5ccbc8', 'check_thunderbird_imap_selection__7fc0ec566fbfa7a7ef758884ad90b95b_qw35sft2_79630b2f', 'check_tb_folder_two_filters__bc08452635875f98955858c473a86873_qw35sft2_861d1f64', 'check_tb_dark_and_mailcheck__9d9a8e81bfc4920e05d53618b0a3d605_qw35sft2_b60e5194', 'check_tb_triple_prefs__b8652f24227f5efb14ee6295b7fbd7cd_qw35sft2_c2a003f2', 'check_thunderbird_nested_folders__1d59e1d1688abeb99fc9fce073e49b55_qw35sft2_181233cc', 'check_tb_acct_smtp__93d97cb63d509700a5843e46760d10d0_qw35sft2_290e5adb', 'check_thunderbird_smtp_full__4c184097a04ef7d98ec36b30ff774f5c_qw35sft2_4d4b6187', 'check_thunderbird_filter_name__e046edd8a7ea67608e0814ad68f68ef0_qw35sft2_e3e859f2', 'check_thunderbird_compose_subject_attachment__59182577b48daf97ea5874cad0371f93_qw35sft2_611d6e51', 'check_tb_dual_filters__f4464353c01b135bac464c93385b0943_qw35sft2_6afd903c', 'check_thunderbird_imap_host__598710d21fbdbbb5b9d697d62aaf5b30_qw35sft2_1bce4b4d', 'check_tb_three_folders_one_filter__9865f0d026c5662c0f08b728c450437f_qw35sft2_c48a3a11', 'check_tb_dual_auto_quote_sig__beb90f2e6aa8d8f432b234d6c0792bb3_qw35sft2_4a8eaccc', 'check_thunderbird_sig_reply__f88a527105940879c254b88e4c3fea68_qw35sft2_75e8712f', 'check_thunderbird_folder_state__50e2868517db7711c169fe6c468e7766_qw35sft2_74469498', 'check_tb_acct_removal__3f1fce882b4fa7de95c0c939ed3f5b9d_qw35sft2_45b0f559', 'check_thunderbird_filter_match_all__f5a0e9992343a59dd61c9a37c78c3db0_qw35sft2_0e8e06e1', 'check_thunderbird_draft_attachment__ac90f87152777ec474b581760325d43f_qw35sft2_b8d2d87e', 'check_tb_incoming_forward__0ab002743507f6ef7d004b0ee40adbf4_qw35sft2_49bcd52a', 'check_thunderbird_smtp_config__efa5fdf04c5026bfc6f8d6ffd452c5ec_qw35sft2_1a67fc58', 'check_thunderbird_imap_port__c8735abb904c2556124a9fb1090745a7_qw35sft2_3d0f67c4', 'check_tb_two_folders_two_filters__85999e7d7f538cb3e1b971a9ebe1ed0f_qw35sft2_a17f0d37', 'check_thunderbird_local_folders__e760cb4cf85af456bb04229f8025e52c_qw35sft2_dbbf05b0', 'check_tb_acct_trash__8fc7e10e457380f7cb927cd320b9664e_qw35sft2_94d91e0d', 'check_thunderbird_filter_and_pref__709b71706c15c6652aa94ca6ff4b6ae6_qw35sft2_7970eb99', 'check_tb_matchall_forward__bd4f2f9ff4f1fb7ec8f18ecd69248d51_qw35sft2_d68f07b8', 'check_thunderbird_name_email__2d78e59033d5901e822a6624aeea3bc7_qw35sft2_9779da52', 'check_thunderbird_compose_cc_attachment__433cf4a73d1a380efac6f447ba251498_qw35sft2_3afa1ea1', 'check_thunderbird_smtp_description__8a8a5b00b6516b6e1171569c19648a1b_qw35sft2_31c7c326', 'check_tb_two_folders_and_filter__951712fea58d5a6835d2f3a50fd315b8_qw35sft2_5acce8a3', 'check_tb_theme_and_folder__bdfb5942d44e1763e7f54cb6f6b53816_qw35sft2_649cfd7b', 'check_tb_dual_auto_quote_reply__07795e258bd2a0b78003c553ab0b53c5_qw35sft2_89499afe', 'check_thunderbird_filter_incoming__4d2f365259c66d4c891bf0c4c8a5a5f0_qw35sft2_3611a8b6', 'check_tb_acct_filter__ddc543b6ee27062796019a5514b0e693_qw35sft2_4ba07b36', 'check_tb_named_forward_filter__e2cb27642e2d718992a0ab3ec3b240f1_qw35sft2_b0ee0ff5', 'check_thunderbird_compose_bcc_attachment__d67afa1e0d57a9bbcae95982d894cc83_qw35sft2_f2e601cb', 'check_thunderbird_smtp_security__9cc78735939880d6658a6582ff470dcd_qw35sft2_f783792d', 'check_tb_two_folders_two_filters_v2__627b4b586c777792c6f3d931dd087820_qw35sft2_16139c00', 'check_vscode_locale_and_minimap__e258fe9bb867a78f6ed6748a6bee2270_qw35sft2_f2661f94', 'check_ext_and_minimap__8a991aa9f9913ccca00f4b1c76aac764_qw35sft2_e803719d']

def check_thunderbird_identity__4a5ace1e770bd482f5a0433cd05f4404(result, expected, **options):
    """Check if Thunderbird identity field matches expected value."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else None
    expected_value = expected.get('expected_value')
    if actual is None or expected_value is None:
        return 0.0
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_thunderbird_smtp_full__2c49114272e52893d5a44b081f8b14e0(result, expected, **options):
    """Check SMTP server configuration including description and username.

    Partial credit:
      - 0.3 for correct hostname
      - 0.2 for correct port
      - 0.2 for correct connection security (try_ssl)
      - 0.3 for correct username
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    smtp_servers = result.get('smtp_servers', {})
    expected_hostname = expected.get('hostname')
    expected_port = expected.get('port')
    expected_ssl = expected.get('try_ssl')
    expected_username = expected.get('username')
    best_score = 0.0
    for (server_id, settings) in smtp_servers.items():
        score = 0.0
        if expected_hostname and settings.get('hostname') == expected_hostname:
            score += 0.3
        if expected_port is not None and settings.get('port') == expected_port:
            score += 0.2
        if expected_ssl is not None and settings.get('try_ssl') == expected_ssl:
            score += 0.2
        if expected_username and settings.get('username') == expected_username:
            score += 0.3
        best_score = max(best_score, score)
    return min(best_score, 1.0)

def check_tb_folder_created__4a5755b15fd8848dbc70d01e3dbb4e89(result, expected, **options):
    """Check if a specific folder exists in the Local Folders listing."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    listing = result.get('listing', '')
    folder_name = expected.get('folder_name', '')
    if not folder_name:
        return 0.0
    lines = [line.strip() for line in listing.strip().split('\n') if line.strip()]
    if folder_name in lines:
        return 1.0
    if folder_name in listing:
        return 1.0
    return 0.0

def check_thunderbird_folder_state__168d8b8e9196711d91fc3322429f68b5(result, expected, **options):
    """Check that specified folders are absent from Thunderbird local folders."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    items = result.get('items', []) if isinstance(result, dict) else []
    expected_absent = expected.get('expected_absent', [])
    if not expected_absent:
        return 0.0
    score = 0.0
    for folder in expected_absent:
        if folder not in items:
            score += 1.0 / len(expected_absent)
    return min(score, 1.0)

def check_tb_msg_count__c7401fcb5b9f30e5bd0dd82a92e3f2df(result, expected, **options):
    """Check that message count in daily folder matches expected count."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_count = result.get('count', -1)
    expected_count = expected.get('expected_count', -1)
    if actual_count == expected_count:
        return 1.0
    return 0.0

def check_tb_email_moved_to_bills__d9f2b609b5f34a856608f932bab06223(result, expected, **options):
    """Check that email was moved from daily to Bills folder. Partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    daily_count = result.get('daily_count', -1)
    bills_has_msg = result.get('bills_has_msg', 0)
    expected_daily_count = expected.get('expected_daily_count', -1)
    if daily_count == expected_daily_count:
        score += 0.5
    if bills_has_msg > 0:
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_account_config__20583839ff43876dc06848262d06ce9a(result, expected, **options):
    """Check email account configuration matches expected values.

    Partial credit:
      - 0.5 for correct email identity
      - 0.3 for correct incoming server hostname
      - 0.2 for correct server type (imap/pop3)
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_email = expected.get('email')
    expected_hostname = expected.get('incoming_hostname')
    expected_type = expected.get('server_type', 'imap')
    identities = result.get('identities', {})
    for (id_name, settings) in identities.items():
        if settings.get('useremail') == expected_email:
            score += 0.5
            break
    servers = result.get('servers', {})
    for (server_name, settings) in servers.items():
        if settings.get('hostname') == expected_hostname:
            score += 0.3
            if settings.get('type') == expected_type:
                score += 0.2
            break
    return min(score, 1.0)

def check_thunderbird_smtp_config__22e52b73b299835b6ec6765d839ed989(result, expected, **options):
    """Check SMTP server configuration matches expected values.

    Partial credit:
      - 0.5 for correct hostname
      - 0.25 for correct port
      - 0.25 for correct connection security (try_ssl)
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    smtp_servers = result.get('smtp_servers', {})
    expected_hostname = expected.get('hostname')
    expected_port = expected.get('port')
    expected_ssl = expected.get('try_ssl')
    best_score = 0.0
    for (server_id, settings) in smtp_servers.items():
        score = 0.0
        if expected_hostname and settings.get('hostname') == expected_hostname:
            score += 0.5
        if expected_port is not None and settings.get('port') == expected_port:
            score += 0.25
        if expected_ssl is not None and settings.get('try_ssl') == expected_ssl:
            score += 0.25
        best_score = max(best_score, score)
    return min(best_score, 1.0)

def check_thunderbird_local_folder__4de43750fd473996ecf007ebf44296c8(result, expected, **options):
    """Check if expected folders exist in Thunderbird local folders."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    items = result.get('items', []) if isinstance(result, dict) else []
    expected_present = expected.get('expected_present', [])
    if not expected_present:
        return 0.0
    score = 0.0
    for folder in expected_present:
        if folder in items:
            score += 1.0 / len(expected_present)
    return min(score, 1.0)

def check_git_email_and_push__0c30b3c3ecef4bb03c9efc8f2c5c6e42_qw35sft2_f4164a4f(result, expected, **options):
    """
    Partial-credit metric checking:
      0.5 — Local git user.email in binder project matches expected value
      0.5 — Remote repo's latest commit message contains expected commit message
    Returns float in [0.0, 1.0].
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_email = expected.get('user_email', 'dev@example.com')
    actual_email = result.get('user_email', '')
    if actual_email == expected_email:
        score += 0.5
    expected_msg = expected.get('commit_message', 'daily update')
    remote_log = result.get('remote_log', '')
    if expected_msg.lower() in remote_log.lower():
        score += 0.5
    return min(score, 1.0)

def check_tb_dark_and_reply__43d8038e02a0c2b1b0f063b3804b00e6_qw35sft2_12b5fb73(result, expected, **options):
    """
    Partial-credit metric: reads local prefs.js path.
      0.5  - extensions.activeThemeID == thunderbird-compact-dark@mozilla.org
      0.5  - mail.identity.id1.reply_on_top == 0  (reply below quote)
    Returns float in [0.0, 1.0].
    """
    import re
    import json
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8', errors='replace') as f:
            prefs_text = f.read()
    except Exception:
        return 0.0
    pref_re = re.compile('^user_pref\\("(?P<key>[^"]+)",\\s*(?P<val>.+)\\);$', re.MULTILINE)
    prefs = {}
    for m in pref_re.finditer(prefs_text):
        key = m.group('key')
        try:
            val = json.loads(m.group('val'))
        except Exception:
            val = m.group('val')
        prefs[key] = val
    score = 0.0
    dark_theme_id = expected.get('dark_theme_id', 'thunderbird-compact-dark@mozilla.org')
    if prefs.get('extensions.activeThemeID') == dark_theme_id:
        score += 0.5
    reply_key = expected.get('reply_key', 'mail.identity.id1.reply_on_top')
    expected_reply_val = expected.get('reply_value', 0)
    actual_reply = prefs.get(reply_key)
    try:
        if int(actual_reply) == int(expected_reply_val):
            score += 0.5
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_thunderbird_sig_name__840f3ea7ded7f1b6695bc630a0a035b7_qw35sft2_00505ae2(result, expected, **options):
    """Check Thunderbird signature text and account display name with partial credit (0.5 each)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_sig = expected.get('expected_sig', '')
    actual_sig = result.get('mail.identity.id1.htmlSigText') or ''
    if expected_sig and expected_sig in actual_sig:
        score += 0.5
    expected_name = expected.get('expected_full_name', '')
    actual_name = result.get('mail.identity.id1.fullName') or ''
    if expected_name and expected_name.strip() == actual_name.strip():
        score += 0.5
    return score

def check_thunderbird_local_folders__1ea15ae7572bc843dca8cde6cbcd3a7c_qw35sft2_eaf2591f(result, expected, **options):
    """
    Partial-credit metric: award equal weight per required folder found.
    expected (after get_rule unwrap): {"required_folders": ["COMPANY", "UNIVERSITY", "RESEARCH"]}
    result: dict mapping folder_name -> bool
    """
    if not isinstance(result, dict):
        return 0.0
    required = expected.get('required_folders', [])
    if not required:
        return 0.0
    score_per = 1.0 / len(required)
    score = 0.0
    for folder in required:
        if result.get(folder, False):
            score += score_per
            logger_qw35sft2_fca153.info('Folder %s: FOUND (+%.3f)', folder, score_per)
        else:
            logger_qw35sft2_fca153.info('Folder %s: MISSING', folder)
    return min(round(score, 4), 1.0)

def check_tb_forward_filter_state__f7fbe032caaa49b6e5ebdeafec739cf8_qw35sft2_66712d03(result, expected, **options):
    """Check that a forward filter exists with the expected destination email.

    Returns 1.0 only if a filter with a forward action exists and the
    actionValue matches the expected forward_to email (case-insensitive).
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count', 0)):
        return 0.0
    expected_email = expected.get('forward_to', '').lower().strip()
    actual_email = (result.get('forward_to') or '').lower().strip()
    has_forward = result.get('has_forward_filter', False)
    if not has_forward:
        return 0.0
    if expected_email and actual_email == expected_email:
        return 1.0
    if not expected_email and has_forward:
        return 1.0
    return 0.0

def check_thunderbird_smtp_no_incoming__f562ce0912b384f5748be41ba253880a_qw35sft2_38a51261(result, expected, **options):
    """Check SMTP configured (0.4) + no incoming accounts (0.4) + correct username (0.2)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_username = expected.get('expected_username', '')
    expected_has_smtp = expected.get('expected_has_smtp', True)
    expected_no_incoming = expected.get('expected_no_incoming', True)
    if expected_has_smtp and result.get('has_smtp'):
        score += 0.4
    if expected_no_incoming and (not result.get('has_incoming')):
        score += 0.4
    if expected_username and result.get('smtp_username', '').lower() == expected_username.lower():
        score += 0.2
    return min(score, 1.0)

def check_tb_acct_folder__2e9b3985bbecaa462dbe43eedf3d22fe_qw35sft2_82ec5e33(result, expected, **options):
    """Partial credit: 0.5 for account removed, 0.5 for Projects folder created."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('account_removed') == expected.get('account_removed', True):
        score += 0.5
    expected_folder = expected.get('folder_name', 'Projects')
    if result.get('folder_exists', False):
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_draft_cc_attachment__ae64d5efe053bb1ed61ef261c2f7eddd_qw35sft2_cb5ccbc8(result, expected, **options):
    """Partial credit: 0.5 for draft with attachment, 0.5 for draft with CC field."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('draft_with_attachment'):
        score += 0.5
    if result.get('draft_with_cc'):
        score += 0.5
    return score

def check_thunderbird_imap_selection__7fc0ec566fbfa7a7ef758884ad90b95b_qw35sft2_79630b2f(result, expected, **options):
    """Check that email is entered and IMAP protocol is visible in Account Setup.
    Partial credit: 0.5 for email present, 0.5 for IMAP visible."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    if result.get('email_present', False):
        score += 0.5
    if result.get('imap_present', False):
        score += 0.5
    return min(score, 1.0)

def check_tb_folder_two_filters__bc08452635875f98955858c473a86873_qw35sft2_861d1f64(result, expected, **options):
    """Score: Promotions folder (0.34) + discount filter (0.33) + sale filter (0.33)."""
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('promotions_exists'):
        score += 0.34
    if result.get('discount_filter'):
        score += 0.33
    if result.get('sale_filter'):
        score += 0.33
    return min(score, 1.0)

def check_tb_dark_and_mailcheck__9d9a8e81bfc4920e05d53618b0a3d605_qw35sft2_b60e5194(result, expected, **options):
    """
    Partial-credit metric: reads local prefs.js path.
      0.5  - extensions.activeThemeID == thunderbird-compact-dark@mozilla.org
      0.5  - mail.server.server1.check_new_mail == True
    Returns float in [0.0, 1.0].
    """
    import re
    import json
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8', errors='replace') as f:
            prefs_text = f.read()
    except Exception:
        return 0.0
    pref_re = re.compile('^user_pref\\("(?P<key>[^"]+)",\\s*(?P<val>.+)\\);$', re.MULTILINE)
    prefs = {}
    for m in pref_re.finditer(prefs_text):
        key = m.group('key')
        try:
            val = json.loads(m.group('val'))
        except Exception:
            val = m.group('val')
        prefs[key] = val
    score = 0.0
    dark_theme_id = expected.get('dark_theme_id', 'thunderbird-compact-dark@mozilla.org')
    if prefs.get('extensions.activeThemeID') == dark_theme_id:
        score += 0.5
    mail_check_key = expected.get('mail_check_key', 'mail.server.server1.check_new_mail')
    if prefs.get(mail_check_key) is True:
        score += 0.5
    return min(score, 1.0)

def check_tb_triple_prefs__b8652f24227f5efb14ee6295b7fbd7cd_qw35sft2_c2a003f2(result, expected, **options):
    """
    Partial-credit metric for three Thunderbird Composition & Addressing goals:
      1. auto_quote disabled  (mail.identity.id1.auto_quote == False)        → 0.34 pts
      2. reply placed below the quote (mail.identity.id1.reply_on_top == 0)  → 0.33 pts
      3. HTML compose disabled (mail.identity.id1.compose_html == False)      → 0.33 pts

    `result` is the raw bytes/string content of prefs.js (from vm_file getter).
    `expected` is the already-unwrapped rules dict.
    """
    if result is None:
        return 0.0
    if isinstance(result, bytes):
        content = result.decode('utf-8', errors='ignore')
    elif isinstance(result, str):
        content = result
    else:
        return 0.0
    prefs = {}
    for line in content.splitlines():
        m = re.match('\\s*user_pref\\("([^"]+)",\\s*(.+)\\);\\s*$', line)
        if m:
            key = m.group(1)
            val_str = m.group(2).strip()
            if val_str == 'true':
                val = True
            elif val_str == 'false':
                val = False
            else:
                try:
                    val = int(val_str)
                except ValueError:
                    try:
                        val = float(val_str)
                    except ValueError:
                        val = val_str.strip('"')
            prefs[key] = val
    score = 0.0
    auto_quote = prefs.get('mail.identity.id1.auto_quote', True)
    if auto_quote is False or auto_quote == False:
        score += 0.34
    reply_on_top = prefs.get('mail.identity.id1.reply_on_top', 1)
    if reply_on_top == 0:
        score += 0.33
    compose_html = prefs.get('mail.identity.id1.compose_html', True)
    if compose_html is False or compose_html == False:
        score += 0.33
    return min(score, 1.0)

def check_thunderbird_nested_folders__1d59e1d1688abeb99fc9fce073e49b55_qw35sft2_181233cc(result, expected, **options):
    """
    Partial-credit metric for 4 checks (0.25 each):
    - COMPANY folder exists
    - UNIVERSITY folder exists
    - INBOX subfolder under COMPANY exists
    - INBOX subfolder under UNIVERSITY exists
    expected (after get_rule unwrap): {
        "company": true, "university": true,
        "company_inbox": true, "university_inbox": true
    }
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    checks = ['company', 'university', 'company_inbox', 'university_inbox']
    weight = 0.25
    for key in checks:
        expected_val = expected.get(key, False)
        actual_val = result.get(key, False)
        if actual_val == expected_val:
            score += weight
            logger_qw35sft2_2fd121.info('Check %s: PASS (+%.2f)', key, weight)
        else:
            logger_qw35sft2_2fd121.info('Check %s: FAIL (got %s, expected %s)', key, actual_val, expected_val)
    return min(round(score, 4), 1.0)

def check_tb_acct_smtp__93d97cb63d509700a5843e46760d10d0_qw35sft2_290e5adb(result, expected, **options):
    """Partial credit: 0.5 for account removed from prefs.js, 0.5 for SMTP server removed."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('account_removed') == expected.get('account_removed', True):
        score += 0.5
    if result.get('smtp_removed') == expected.get('smtp_removed', True):
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_smtp_full__4c184097a04ef7d98ec36b30ff774f5c_qw35sft2_4d4b6187(result, expected, **options):
    """Check all 4 SMTP fields with partial credit: 0.25 each for hostname, port, try_ssl, username."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_hostname = expected.get('expected_hostname', '')
    expected_port = expected.get('expected_port', 0)
    expected_try_ssl = expected.get('expected_try_ssl', -1)
    expected_username = expected.get('expected_username', '')
    if expected_hostname and result.get('hostname', '').lower() == expected_hostname.lower():
        score += 0.25
    if expected_port and result.get('port', 0) == expected_port:
        score += 0.25
    if expected_try_ssl >= 0 and result.get('try_ssl', -1) == expected_try_ssl:
        score += 0.25
    if expected_username and result.get('username', '').lower() == expected_username.lower():
        score += 0.25
    return min(score, 1.0)

def check_thunderbird_filter_name__e046edd8a7ea67608e0814ad68f68ef0_qw35sft2_e3e859f2(result, expected, **options):
    """Check that a filter with the expected name exists."""
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_names')):
        return 0.0
    expected_name = expected.get('filter_name', '')
    filter_names = result.get('filter_names', [])
    if expected_name and any((expected_name.lower() in n.lower() for n in filter_names)):
        return 1.0
    return 0.0

def check_thunderbird_compose_subject_attachment__59182577b48daf97ea5874cad0371f93_qw35sft2_611d6e51(result, expected, **options):
    """Partial credit: 0.5 for attachment present, 0.5 for updated subject present."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_attachment'):
        score += 0.5
    if result.get('has_subject'):
        score += 0.5
    return score

def check_tb_dual_filters__f4464353c01b135bac464c93385b0943_qw35sft2_6afd903c(result, expected, **options):
    """Check forward filter (0.5) and mark-as-read filter (0.5).

    Partial credit:
    - 0.5 if a forward filter exists targeting the expected email
    - 0.5 if a mark-as-read filter exists
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count', 0)):
        return 0.0
    score = 0.0
    expected_email = expected.get('forward_to', '').lower().strip()
    actual_email = (result.get('forward_to') or '').lower().strip()
    if result.get('has_forward', False):
        if not expected_email or actual_email == expected_email:
            score += 0.5
    if result.get('has_mark_read', False):
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_imap_host__598710d21fbdbbb5b9d697d62aaf5b30_qw35sft2_1bce4b4d(result, expected, **options):
    """Check email entered and IMAP hostname (outlook.office365.com) is visible in manual config.
    Partial credit: 0.4 for email, 0.6 for IMAP hostname in manual configuration view."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    if result.get('email_present', False):
        score += 0.4
    if result.get('imap_host_present', False):
        score += 0.6
    return min(score, 1.0)

def check_tb_three_folders_one_filter__9865f0d026c5662c0f08b728c450437f_qw35sft2_c48a3a11(result, expected, **options):
    """Score: Promotions folder (0.25) + Deals folder (0.25) + Archive folder (0.25) + discount filter (0.25)."""
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('promotions_exists'):
        score += 0.25
    if result.get('deals_exists'):
        score += 0.25
    if result.get('archive_exists'):
        score += 0.25
    if result.get('discount_filter'):
        score += 0.25
    return min(score, 1.0)

def check_tb_dual_auto_quote_sig__beb90f2e6aa8d8f432b234d6c0792bb3_qw35sft2_4a8eaccc(result, expected, **options):
    """
    Partial-credit metric for two Thunderbird Composition & Addressing goals:
      1. auto_quote disabled  (mail.identity.id1.auto_quote == False)     → 0.5 pts
      2. signature disabled in replies (mail.identity.id1.sig_on_reply == False) → 0.5 pts

    `result` is the raw bytes/string content of prefs.js (from vm_file getter).
    `expected` is the already-unwrapped rules dict.
    """
    if result is None:
        return 0.0
    if isinstance(result, bytes):
        content = result.decode('utf-8', errors='ignore')
    elif isinstance(result, str):
        content = result
    else:
        return 0.0
    prefs = {}
    for line in content.splitlines():
        m = re.match('\\s*user_pref\\("([^"]+)",\\s*(.+)\\);\\s*$', line)
        if m:
            key = m.group(1)
            val_str = m.group(2).strip()
            if val_str == 'true':
                val = True
            elif val_str == 'false':
                val = False
            else:
                try:
                    val = int(val_str)
                except ValueError:
                    try:
                        val = float(val_str)
                    except ValueError:
                        val = val_str.strip('"')
            prefs[key] = val
    score = 0.0
    auto_quote = prefs.get('mail.identity.id1.auto_quote', True)
    if auto_quote is False or auto_quote == False:
        score += 0.5
    sig_on_reply = prefs.get('mail.identity.id1.sig_on_reply', True)
    if sig_on_reply is False or sig_on_reply == False:
        score += 0.5
    return score

def check_thunderbird_sig_reply__f88a527105940879c254b88e4c3fea68_qw35sft2_75e8712f(result, expected, **options):
    """Check Thunderbird signature text and reply-to address with partial credit (0.5 each)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_sig = expected.get('expected_sig', '')
    actual_sig = result.get('mail.identity.id1.htmlSigText') or ''
    if expected_sig and expected_sig in actual_sig:
        score += 0.5
    expected_reply = expected.get('expected_reply_to', '')
    actual_reply = result.get('mail.identity.id1.replyTo') or ''
    if expected_reply and expected_reply.strip().lower() == actual_reply.strip().lower():
        score += 0.5
    return score

def check_thunderbird_folder_state__50e2868517db7711c169fe6c468e7766_qw35sft2_74469498(result, expected, **options):
    """
    Partial-credit metric for COMPANY (0.34) + UNIVERSITY (0.33) + PROJECTS-in-COMPANY (0.33).
    expected (after get_rule unwrap): {"company": true, "university": true, "projects_in_company": true}
    result: dict with boolean values from getter.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    weights = {'company': 0.34, 'university': 0.33, 'projects_in_company': 0.33}
    for key, weight in weights.items():
        expected_val = expected.get(key, False)
        actual_val = result.get(key, False)
        if actual_val == expected_val:
            score += weight
            logger_qw35sft2_34eb84.info('Check %s: PASS (+%.2f)', key, weight)
        else:
            logger_qw35sft2_34eb84.info('Check %s: FAIL (got %s, expected %s)', key, actual_val, expected_val)
    return min(round(score, 4), 1.0)

def check_tb_acct_removal__3f1fce882b4fa7de95c0c939ed3f5b9d_qw35sft2_45b0f559(result, expected, **options):
    """Return 1.0 if the account email is absent from prefs.js, else 0.0."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_removed = expected.get('account_removed', True)
    actual_removed = result.get('account_removed', False)
    return 1.0 if actual_removed == expected_removed else 0.0

def check_thunderbird_filter_match_all__f5a0e9992343a59dd61c9a37c78c3db0_qw35sft2_0e8e06e1(result, expected, **options):
    """Check that at least one filter uses 'Match all messages' (condition=ALL)."""
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count')):
        return 0.0
    expected_val = expected.get('has_match_all_filter', True)
    actual_val = result.get('has_match_all_filter', False)
    return 1.0 if actual_val == expected_val else 0.0

def check_thunderbird_draft_attachment__ac90f87152777ec474b581760325d43f_qw35sft2_b8d2d87e(result, expected, **options):
    """Return 1.0 if Thunderbird Drafts contains the aws-bill.pdf attachment."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('draft_with_attachment') else 0.0

def check_tb_incoming_forward__0ab002743507f6ef7d004b0ee40adbf4_qw35sft2_49bcd52a(result, expected, **options):
    """Check Getting New Mail trigger with forward action (0.5) and destination email (0.5).

    Partial credit:
    - 0.5 if a filter exists with Getting New Mail (type bit 0x1) and a forward action
    - 0.5 if the forward destination email matches expected forward_to (case-insensitive)
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count', 0)):
        return 0.0
    score = 0.0
    if result.get('has_incoming_forward', False):
        score += 0.5
    expected_email = expected.get('forward_to', '').lower().strip()
    actual_email = (result.get('forward_to') or '').lower().strip()
    if expected_email and actual_email == expected_email:
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_smtp_config__efa5fdf04c5026bfc6f8d6ffd452c5ec_qw35sft2_1a67fc58(result, expected, **options):
    """Check Thunderbird SMTP hostname and username. 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_hostname = expected.get('expected_hostname', '')
    expected_username = expected.get('expected_username', '')
    actual_hostname = result.get('hostname', '')
    actual_username = result.get('username', '')
    if expected_hostname and actual_hostname.lower() == expected_hostname.lower():
        score += 0.5
    if expected_username and actual_username.lower() == expected_username.lower():
        score += 0.5
    return score

def check_thunderbird_imap_port__c8735abb904c2556124a9fb1090745a7_qw35sft2_3d0f67c4(result, expected, **options):
    """Check email, IMAP hostname (outlook.office365.com), and port 993 in manual config.
    Partial credit: 0.33 for email, 0.34 for IMAP hostname, 0.33 for port 993."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    if result.get('email_present', False):
        score += 0.33
    if result.get('imap_host_present', False):
        score += 0.34
    if result.get('port_993_present', False):
        score += 0.33
    return min(score, 1.0)

def check_tb_two_folders_two_filters__85999e7d7f538cb3e1b971a9ebe1ed0f_qw35sft2_a17f0d37(result, expected, **options):
    """Score: Promotions folder (0.25) + Newsletter folder (0.25) + discount filter (0.25) + newsletter filter (0.25)."""
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('promotions_exists'):
        score += 0.25
    if result.get('newsletter_exists'):
        score += 0.25
    if result.get('discount_filter'):
        score += 0.25
    if result.get('newsletter_filter'):
        score += 0.25
    return min(score, 1.0)

def check_thunderbird_local_folders__e760cb4cf85af456bb04229f8025e52c_qw35sft2_dbbf05b0(result, expected, **options):
    """
    Partial-credit metric: award equal weight per required folder found.
    expected (after get_rule unwrap): {"required_folders": ["COMPANY", "UNIVERSITY", "WORK", "PERSONAL"]}
    result: dict mapping folder_name -> bool
    """
    if not isinstance(result, dict):
        return 0.0
    required = expected.get('required_folders', [])
    if not required:
        return 0.0
    score_per = 1.0 / len(required)
    score = 0.0
    for folder in required:
        if result.get(folder, False):
            score += score_per
            logger_qw35sft2_1d640f.info('Folder %s: FOUND (+%.3f)', folder, score_per)
        else:
            logger_qw35sft2_1d640f.info('Folder %s: MISSING', folder)
    return min(round(score, 4), 1.0)

def check_tb_acct_trash__8fc7e10e457380f7cb927cd320b9664e_qw35sft2_94d91e0d(result, expected, **options):
    """Partial credit: 0.5 for account removed, 0.5 for Empty Trash on Exit enabled."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('account_removed') == expected.get('account_removed', True):
        score += 0.5
    if result.get('empty_trash_enabled') == expected.get('empty_trash_enabled', True):
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_filter_and_pref__709b71706c15c6652aa94ca6ff4b6ae6_qw35sft2_7970eb99(result, expected, **options):
    """
    Partial credit check:
    - 0.5 for enabling applyIncomingFilters preference
    - 0.5 for creating at least one message filter
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('apply_incoming_filters') == expected.get('apply_incoming_filters', True):
        score += 0.5
    if result.get('has_filter') == expected.get('has_filter', True):
        score += 0.5
    return min(score, 1.0)

def check_tb_matchall_forward__bd4f2f9ff4f1fb7ec8f18ecd69248d51_qw35sft2_d68f07b8(result, expected, **options):
    """Check match-all condition (0.5) and forward destination email (0.5).

    Partial credit:
    - 0.5 if at least one filter uses condition=ALL (match all messages)
    - 0.5 if a forward action targets the expected email address
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count', 0)):
        return 0.0
    score = 0.0
    if result.get('has_match_all', False) == expected.get('has_match_all', True):
        score += 0.5
    expected_email = expected.get('forward_to', '').lower().strip()
    actual_email = (result.get('forward_to') or '').lower().strip()
    if expected_email and actual_email == expected_email:
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_name_email__2d78e59033d5901e822a6624aeea3bc7_qw35sft2_9779da52(result, expected, **options):
    """Check that the display name 'Anonym X' and email are both present in Account Setup.
    Partial credit: 0.5 for name, 0.5 for email."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    if result.get('name_present', False):
        score += 0.5
    if result.get('email_present', False):
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_compose_cc_attachment__433cf4a73d1a380efac6f447ba251498_qw35sft2_3afa1ea1(result, expected, **options):
    """Partial credit: 0.5 for attachment present, 0.5 for CC address present."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_attachment'):
        score += 0.5
    if result.get('has_cc'):
        score += 0.5
    return score

def check_thunderbird_smtp_description__8a8a5b00b6516b6e1171569c19648a1b_qw35sft2_31c7c326(result, expected, **options):
    """Check Thunderbird SMTP description. 0.5 for SMTP present, 0.5 for correct description."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_desc_contains = expected.get('expected_description_contains', '')
    if result.get('has_smtp'):
        score += 0.5
    actual_desc = result.get('description', '')
    if expected_desc_contains and expected_desc_contains.lower() in actual_desc.lower():
        score += 0.5
    return score

def check_tb_two_folders_and_filter__951712fea58d5a6835d2f3a50fd315b8_qw35sft2_5acce8a3(result, expected, **options):
    """Score: Promotions folder (0.33) + Deals folder (0.33) + discount filter (0.34)."""
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('promotions_exists'):
        score += 0.33
    if result.get('deals_exists'):
        score += 0.33
    if result.get('discount_filter'):
        score += 0.34
    return min(score, 1.0)

def check_tb_theme_and_folder__bdfb5942d44e1763e7f54cb6f6b53816_qw35sft2_649cfd7b(result, expected, **options):
    """
    Partial-credit metric for Dark theme + Night Notes folder creation.
      0.5  - extensions.activeThemeID == thunderbird-compact-dark@mozilla.org
      0.5  - Night Notes local folder exists (msf file present)
    Returns float in [0.0, 1.0].
    """
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    dark_theme_id = expected.get('dark_theme_id', 'thunderbird-compact-dark@mozilla.org')
    if result.get('active_theme_id') == dark_theme_id:
        score += 0.5
    if result.get('night_notes_folder_exists'):
        score += 0.5
    return min(score, 1.0)

def check_tb_dual_auto_quote_reply__07795e258bd2a0b78003c553ab0b53c5_qw35sft2_89499afe(result, expected, **options):
    """
    Partial-credit metric for two Thunderbird Composition & Addressing goals:
      1. auto_quote disabled  (mail.identity.id1.auto_quote == False)   -> 0.5 pts
      2. reply placed below the quote (mail.identity.id1.reply_on_top == 0) -> 0.5 pts

    `result` is the raw bytes/string content of prefs.js (from vm_file getter).
    `expected` is the already-unwrapped rules dict (framework strips 'rules' wrapper).
    """
    if result is None:
        return 0.0
    if isinstance(result, bytes):
        content = result.decode('utf-8', errors='ignore')
    elif isinstance(result, str):
        content = result
    else:
        return 0.0
    prefs = {}
    for line in content.splitlines():
        m = re.match('\\s*user_pref\\("([^"]+)",\\s*(.+)\\);\\s*$', line)
        if m:
            key = m.group(1)
            val_str = m.group(2).strip()
            if val_str == 'true':
                val = True
            elif val_str == 'false':
                val = False
            else:
                try:
                    val = int(val_str)
                except ValueError:
                    try:
                        val = float(val_str)
                    except ValueError:
                        val = val_str.strip('"')
            prefs[key] = val
    score = 0.0
    auto_quote = prefs.get('mail.identity.id1.auto_quote', True)
    if auto_quote is False:
        score += 0.5
    reply_on_top = prefs.get('mail.identity.id1.reply_on_top', 1)
    if reply_on_top == 0:
        score += 0.5
    return score

def check_thunderbird_filter_incoming__4d2f365259c66d4c891bf0c4c8a5a5f0_qw35sft2_3611a8b6(result, expected, **options):
    """Check that at least one filter has the Getting New Mail (Incoming) trigger enabled."""
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count')):
        return 0.0
    expected_val = expected.get('has_incoming_filter', True)
    actual_val = result.get('has_incoming_filter', False)
    return 1.0 if actual_val == expected_val else 0.0

def check_tb_acct_filter__ddc543b6ee27062796019a5514b0e693_qw35sft2_4ba07b36(result, expected, **options):
    """Partial credit: 0.5 for account removed, 0.5 for newsletter filter created."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('account_removed') == expected.get('account_removed', True):
        score += 0.5
    if result.get('newsletter_filter_exists', False):
        score += 0.5
    return min(score, 1.0)

def check_tb_named_forward_filter__e2cb27642e2d718992a0ab3ec3b240f1_qw35sft2_b0ee0ff5(result, expected, **options):
    """Check filter name (0.5) and forward destination email (0.5).

    Partial credit:
    - 0.5 if any filter name contains the expected filter_name substring (case-insensitive)
    - 0.5 if the forward destination email matches expected forward_to (case-insensitive)
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('filter_count', 0)):
        return 0.0
    score = 0.0
    expected_name = expected.get('filter_name', '').lower().strip()
    filter_names = [n.lower() for n in result.get('filter_names', [])]
    if expected_name and any((expected_name in n for n in filter_names)):
        score += 0.5
    expected_email = expected.get('forward_to', '').lower().strip()
    actual_email = (result.get('forward_to') or '').lower().strip()
    if expected_email and actual_email == expected_email:
        score += 0.5
    return min(score, 1.0)

def check_thunderbird_compose_bcc_attachment__d67afa1e0d57a9bbcae95982d894cc83_qw35sft2_f2e601cb(result, expected, **options):
    """Partial credit: 0.5 for attachment present, 0.5 for BCC address present."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_attachment'):
        score += 0.5
    if result.get('has_bcc'):
        score += 0.5
    return score

def check_thunderbird_smtp_security__9cc78735939880d6658a6582ff470dcd_qw35sft2_f783792d(result, expected, **options):
    """Check Thunderbird SMTP port and connection security. 0.5 each.
    try_ssl: 0=None, 2=STARTTLS, 3=SSL/TLS
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_port = expected.get('expected_port', 0)
    expected_try_ssl = expected.get('expected_try_ssl', -1)
    actual_port = result.get('port', 0)
    actual_try_ssl = result.get('try_ssl', -1)
    if expected_port and actual_port == expected_port:
        score += 0.5
    if expected_try_ssl >= 0 and actual_try_ssl == expected_try_ssl:
        score += 0.5
    return score

def check_tb_two_folders_two_filters_v2__627b4b586c777792c6f3d931dd087820_qw35sft2_16139c00(result, expected, **options):
    """Score: Promotions folder (0.25) + Coupons folder (0.25) + discount filter (0.25) + coupon filter (0.25)."""
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('promotions_exists'):
        score += 0.25
    if result.get('coupons_exists'):
        score += 0.25
    if result.get('discount_filter'):
        score += 0.25
    if result.get('coupon_filter'):
        score += 0.25
    return min(score, 1.0)

def check_vscode_locale_and_minimap__e258fe9bb867a78f6ed6748a6bee2270_qw35sft2_f2661f94(result, expected, **options):
    """Check VS Code display language is German and editor.minimap.enabled is False.

    result: dict with 'locale' and 'settings' from getter
    expected: rules dict (already unwrapped), contains 'locale' and 'minimap_enabled'
    Partial credit: 0.5 per sub-goal.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_locale = expected.get('locale', '').lower()
    actual_locale = result.get('locale', '').lower()
    if expected_locale and (actual_locale == expected_locale or actual_locale.startswith(expected_locale)):
        score += 0.5
    expected_minimap = expected.get('minimap_enabled')
    if expected_minimap is not None:
        settings = result.get('settings', {})
        actual_minimap = settings.get('editor.minimap.enabled', settings.get('editor.minimap', {}).get('enabled', None) if isinstance(settings.get('editor.minimap'), dict) else None)
        if actual_minimap is not None and bool(actual_minimap) == bool(expected_minimap):
            score += 0.5
    return round(score, 2)

def check_ext_and_minimap__8a991aa9f9913ccca00f4b1c76aac764_qw35sft2_e803719d(result, expected, **options):
    """Check extension installed (0.5) + minimap disabled (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'undefined_publisher.test')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    expected_minimap = expected.get('minimap_enabled')
    actual_minimap = result.get('minimap_enabled')
    if expected_minimap is not None and actual_minimap == expected_minimap:
        score += 0.5
    return score
