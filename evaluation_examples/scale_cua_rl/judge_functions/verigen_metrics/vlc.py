"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_vlc_play_and_pause__836e80a44129f44b6386ea1f41fd189b', 'check_video_nosub__bf3c684d8dd921786e84a6a366f46d5a', 'check_vlc_systray__fba3e4f74122f0f091c75248a9112c9b', 'check_vlc_paused__9f6d71727e5c15543b18896fb310844a', 'check_vlc_snapshot_path__5e0183d902d6d20115487cceb9d1bd09', 'check_vlc_snapshot_subdir__5366805ab8e670a4853b7e253af708d2', 'check_vlc_one_instance__e3afb507daace17955585b556bc80462', 'check_vlc_minimal_view__0dfff08b7a8ec4f9e7e69d81a86736ba', 'check_vlc_max_volume__ccc1da9663fce350237c6ebebb163164', 'check_qt_pause_minimized__084cc353323ebeb25e57c8f66bfbb2a1', 'check_vlc_multi_settings__be291c8b2f31ab750f3a52d157cbc9fd', 'check_vlc_play_and_stop__23bd9b1b573f33533486832ddf9f8705', 'check_audio_extraction__e0d33b8eb904792c666a0fe3a33dc4f5', 'check_audio_format__3c3196d8d3d24397883e7bbcfddc8dd0', 'check_vlc_snapshot_file__da77e4abbc92d5a7a2f84a9516d6d0b6', 'check_audio_format__9e6e6e9addabc1fa910c5d7bd33b2f51', 'check_video_dimensions__92a1bbba6e28a320b340946ad5c7a75d', 'check_vlc_dual_snapshot__1118aad2dd728a496771a4813a51c0c2', 'check_audio_clip_format__7c8156c2a09b3bba70e2cfa0aa1a4dcc', 'check_video_dimensions__22c3b2acb75b2231bd0ecc7a3dd0e178', 'check_wav_file_exists__1606287de635101a5a1b8cd46a5bf405', 'check_vlc_snapshot_exists__786ac827ab684b501ddc668a08ec2cfa', 'check_vlc_snapshot_exists__441247156205cdab6fb75baec6278a39', 'check_vlc_network_cache__72b69002c9e414b44114ece1b158222e', 'check_vlc_qt_continue__e543a11713a93c5fff74c1d7c5501303', 'check_video_duration__38add8a045487bc69861547adf883b78_qw35sft2_3c8ca599', 'check_video_duration__dd6b82d69d19c741223e8dfd4c176567_qw35sft2_d3a22df1', 'check_single_mp3_tags__85f79002782d5324b058fe59de70692f_qw35sft2_8c651e51', 'check_single_mp3_tags__367848f954ef0bdc8124c00e1a0efa09_qw35sft2_7fc7746c', 'check_vlc_fullscreen_config__d319a77e7f319beb729da2938a21127c_qw35sft2_875bb6eb', 'check_single_mp3_tags__d208efee5ee527e892400da003674774_qw35sft2_41c19618', 'check_mp3_exists_mp4_removed__e17edfd562e8913a3a72a3bcb87c9c97_qw35sft2_095d9ab3', 'check_vlc_record_path__33b97504f4caa5a543c607fe8892e403_qw35sft2_291ff20a', 'check_vlc_max_volume__95b68fd0ad40f238d037990e09029619_qw35sft2_a93a4cf8', 'check_vlc_dual_config__de221324409107598f868e822754ff7f_qw35sft2_fa346a76', 'check_vlc_snap_and_rename__25f0d8c371fba607f6a6c0535521c207_qw35sft2_f56f74de', 'check_vlc_play_and_random__cd6a6657f06ad537781a5ce037c2b82a_qw35sft2_ab28cbbc', 'check_vlc_brightness__1e9b265786c948c671445fe3130d0b36_qw35sft2_1069eca8', 'check_vlc_stream_and_cache__4f78681b35e8e22495902169ef7b5c87_qw35sft2_680c81a3', 'check_snapshot_and_wallpaper__1e243c15dcedcf4e9a4cacc616a0598b_qw35sft2_a4f44df9', 'check_vlc_recfolder_and_maxvol__517e4427a675aa3407c571ed7524794f_qw35sft2_d7155b18', 'check_vlc_record_path__e488a41c13481ad2b017a00cb14067ec_qw35sft2_8dd20fab', 'check_vlc_max_volume__15cb17bd2e24428fe6d640ea8d5d556a_qw35sft2_af4230f7', 'check_vlc_instance_settings__71c2e6dc1ced259311cab9f07521fd96_qw35sft2_65756d41', 'check_mp3_validity__35ef4e9e6e88fc273aff9d3f203be720_qw35sft2_54cf8bf4', 'check_vlc_play_and_rate__a87ab91c6345361c3c6e374fc52958bf_qw35sft2_949dbd4a', 'check_vlc_rename__bd8f719f06ca22a9a2d291dd693db756_qw35sft2_53a730de', 'check_vlc_cone_and_oneinstance__6535f71afc4db044467cc2ad41010df9_qw35sft2_582f829f', 'check_wallpaper_and_desktop_snapshot__5688825554c17cb3fb1e24246b369618_qw35sft2_53323222', 'check_vlc_stream_and_cone__e0203e1b79d3816d60eca3abd8c45b33_qw35sft2_6d2cf461', 'check_vlc_snap_in_captures__6d219cf2312d28ed73dcf8a9fc508b4d_qw35sft2_288b20f7', 'check_vlc_fullscreen_and_snapshot__b80790f2f04245c0798c2ba24e3ceaf2_qw35sft2_5c114c55', 'check_vlc_maxvol_and_bgcone__f99d034b7e5202fd562c023da2142c74_qw35sft2_66853c69', 'check_vlc_playback_prefs__db101724f343bdcc09eb9f11a0e77c94_qw35sft2_55f6826f', 'check_vlc_effects_dialog__80d9e76aa584c867ce735bfda2f705a1_qw35sft2_b960a33e', 'check_vlc_play_and_repeat__dc3aa9f060ebd375fbf717e1b7e05e93_qw35sft2_b1cfa8c7', 'check_mp3_mp4_both_exist__41a079bbfac3c55bed03337726ec0862_qw35sft2_251d7a95', 'check_vlc_prefs__8b6f202a2cac4c2b2f141e2f0645d358_qw35sft2_e766877d', 'check_vlc_saturation__493da9b8d1a20a26a98a5cc60b751e43_qw35sft2_fdfb9d9d', 'check_snapshot_file_exists__4e0abb3a76e2d3729689698c9288f79c_qw35sft2_402eb8d2', 'check_vlc_cone_and_volume__78f02ce8f924783e3007881645e6184a_qw35sft2_8ddca2e6', 'check_vlc_stream_and_snapshot__b5aec674ad8f49621ab70bae692967c5_qw35sft2_331ca700', 'vlc_snapshot_dual_loc__b8a50137decabc772595757ced9c456a_qw35sft2_d6567a65', 'check_vlc_fullscreen_and_maxvol__ecd5fecf55729059a58a10d83c8c9eeb_qw35sft2_de69a72b', 'check_vlc_triple_prefs__29cde5ba2a87fed510792760a2c64d1f_qw35sft2_b89ee98a', 'check_vlc_play_and_record_path__d9c5d9aac51555091a28197b6da34259_qw35sft2_7910036a', 'check_mp3_file_size__55ceccfa8167b4a354049a9775b28527_qw35sft2_098ff36e', 'check_vlc_maxvol_and_bgcone_expands__d25099f13e7da2513cf33c07be0e7955_qw35sft2_378e65f4', 'check_vlc_threshold__c451d2a7e7a5b8e510cce0c34da4325d_qw35sft2_19eb83a8', 'check_vlc_cone_and_minview__24dc6e88ff00c69857d313c4738a8f31_qw35sft2_4532da52', 'check_vlc_snapshot_wallpaper__d24bdc521260cf2eb34948f145dea565_qw35sft2_66214d9a', 'check_vlc_record_path__f613292e2de2a0e0145aaa24fdebbc72_qw35sft2_67a189da', 'check_vlc_stream_and_volume__a51ef1127f9b7bededbd54fd205d555d_qw35sft2_6a01a4f9', 'vlc_snapshot_pictures__cdcbbd90330cd55be87f2b9353b8c43c_qw35sft2_1ff63cb9', 'check_vlc_fullscreen_and_bgcone__986ee683b9dd7fb42919cda3330ee515_qw35sft2_5e67ef02', 'check_vlc_loop_prefs__a02aa930219f8e3a3317ea91cceaa2ff_qw35sft2_b6c8b6bc', 'check_vlc_play_and_loop__c8546449445cd5e2551a333da05cf053_qw35sft2_2ed6f216', 'check_vlc_contrast__ab50e89f0aef704ab57d4c9d2579417f_qw35sft2_eac2f79d', 'check_vlc_max_volume__a9382a8b8a5c0cd01bfb6f50454b9a65_qw35sft2_e32b8685', 'check_vlc_cone_and_recordpath__914a6c12022e434d9eb3c729da97ace4_qw35sft2_5861b639', 'check_mp3_existence__77b822a22aadd202cc6c547979661a88_qw35sft2_2317004e', 'check_vlc_snap_on_desktop__dc35efec6c78cc3245c9e38e3e5c5d4a_qw35sft2_f8227f2f', 'check_wallpaper_and_vlc_running__2cdc858f118b65a32b6755071b250267_qw35sft2_74c7dcea', 'check_vlc_fullscreen_and_recfolder__27e3d8a3f9225f7e9bc7da9d0f1debd6_qw35sft2_8442bba5', 'check_vlc_stream_and_ontop__5d8c2e81cb9645b71051b528d54492c6_qw35sft2_47d573db', 'check_vlc_prefs__a7d5b4b88d0304ec4fb07fa10130fbc0_qw35sft2_0b6de3eb']

def check_vlc_play_and_pause__836e80a44129f44b6386ea1f41fd189b(actual_config_path, rule):
    """
    Checks if VLC's 'Pause on the last frame of a video' setting (play-and-pause) is enabled.
    Values: 0=disabled (default), 1=enabled
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_play_and_pause = rule.get('expected_play_and_pause', '1')
    if isinstance(expected_play_and_pause, int):
        expected_play_and_pause = str(expected_play_and_pause)
    try:
        play_and_pause = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'play-and-pause=' in line:
                play_and_pause = line.split('=')[-1].strip()
        if play_and_pause == expected_play_and_pause:
            return 1.0
        else:
            logger.warning(f'play-and-pause mismatch: expected={expected_play_and_pause}, actual={play_and_pause}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_nosub__bf3c684d8dd921786e84a6a366f46d5a(result, expected, **options):
    """Check video without subtitles with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('file_exists'):
        score += 0.25
    if result.get('has_video'):
        score += 0.25
    if result.get('has_audio'):
        score += 0.25
    if result.get('no_subtitles'):
        score += 0.25
    return min(score, 1.0)

def check_vlc_systray__fba3e4f74122f0f091c75248a9112c9b(actual_config_path, rule):
    """Check if VLC systray icon setting matches expected value."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_value = rule.get('expected_qt_system_tray')
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    try:
        actual_value = '1'
        for line in config_file.split('\n'):
            if 'qt-system-tray=' in line:
                actual_value = line.split('=')[-1].strip()
        if actual_value == expected_value:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_paused__9f6d71727e5c15543b18896fb310844a(actual_status_path, expected, **options):
    """
    Checks if VLC has a specific file loaded and is in paused state.
    Partial credit: 0.5 for correct file loaded, 0.5 for paused state.
    """
    if not actual_status_path or not os.path.exists(actual_status_path):
        return 0.0
    with open(actual_status_path, 'rb') as file:
        actual_status = file.read().decode('utf-8')
    try:
        tree = ElementTree.fromstring(actual_status)
    except ElementTree.ParseError:
        logger.error('Failed to parse VLC status XML')
        return 0.0
    score = 0.0
    expected_filename = expected.get('file_name', '')
    file_paths = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]']
    file_info = None
    for path in file_paths:
        element = tree.find(path)
        if element is not None and element.text:
            file_info = element.text
            break
    if file_info:
        actual_basename = os.path.basename(file_info)
        if actual_basename == expected_filename or expected_filename in file_info:
            score += 0.5
    state_elem = tree.find('state')
    if state_elem is not None:
        state = state_elem.text
        if state == 'paused':
            score += 0.5
    return score

def check_vlc_snapshot_path__5e0183d902d6d20115487cceb9d1bd09(actual_config_path, rule, **options):
    """Check if VLC's snapshot path is set to the expected value."""
    try:
        with open(actual_config_path, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception:
        return 0.0
    expected_path = rule.get('expected_snapshot_path', '')
    for line in config_file.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'snapshot-path' in line and 'snapshot-path-hierarchical' not in line:
            current_path = line.split('=', 1)[-1].strip()
            if current_path == expected_path:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_vlc_snapshot_subdir__5366805ab8e670a4853b7e253af708d2(result, expected, **options):
    """Check if snapshot exists in the correct subdirectory. Partial credit."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('dir_exists', False):
        score += 0.2
    if result.get('file_exists', False):
        score += 0.4
        mime_type = result.get('mime_type', '')
        if mime_type.startswith('image/'):
            score += 0.2
        file_size = result.get('file_size', 0)
        if file_size > 1024:
            score += 0.2
    return min(score, 1.0)

def check_vlc_one_instance__e3afb507daace17955585b556bc80462(actual_config_path, rule):
    """Check if VLC 'Allow only one instance' setting matches expected value."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_value = rule.get('expected_one_instance')
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    try:
        actual_value = '0'
        for line in config_file.split('\n'):
            if 'one-instance=' in line and 'one-instance-when-started-from-file' not in line:
                actual_value = line.split('=')[-1].strip()
        if actual_value == expected_value:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_minimal_view__0dfff08b7a8ec4f9e7e69d81a86736ba(actual_config_path, rule):
    """Check if VLC 'Start in minimal view mode' setting matches expected value."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_value = rule.get('expected_qt_minimal_view')
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    try:
        actual_value = '0'
        for line in config_file.split('\n'):
            if 'qt-minimal-view=' in line:
                actual_value = line.split('=')[-1].strip()
        if actual_value == expected_value:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_max_volume__ccc1da9663fce350237c6ebebb163164(result, expected, **options):
    """Check if VLC maximum volume is set to the expected value.
    Partial credit:
      0.5 - Setting exists and is uncommented (actively set)
      1.0 - Setting value matches expected
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_value = expected.get('expected_max_volume', 200)
    actual_value = result.get('value')
    is_commented = result.get('is_commented', True)
    if actual_value is None:
        return 0.0
    score = 0.0
    if not is_commented:
        score = 0.5
    if actual_value == expected_value and (not is_commented):
        score = 1.0
    return score

def check_qt_pause_minimized__084cc353323ebeb25e57c8f66bfbb2a1(actual_config_path, rule):
    """Check if VLC's 'Pause playback when minimized' setting matches expected value."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected = rule.get('expected_qt_pause_minimized', '1')
    if isinstance(expected, int):
        expected = str(expected)
    try:
        qt_pause_minimized = '0'
        for line in config_file.split('\n'):
            if 'qt-pause-minimized=' in line:
                qt_pause_minimized = line.split('=')[-1].strip()
        if qt_pause_minimized == expected:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_multi_settings__be291c8b2f31ab750f3a52d157cbc9fd(result, expected, **options):
    """Check multiple VLC Qt settings with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_max_volume = expected.get('expected_qt_max_volume', '200')
    actual_max_volume = result.get('qt_max_volume', '125')
    if str(actual_max_volume) == str(expected_max_volume):
        score += 0.5
    expected_bgcone = expected.get('expected_qt_bgcone', '0')
    actual_bgcone = result.get('qt_bgcone', '1')
    if str(actual_bgcone) == str(expected_bgcone):
        score += 0.5
    return min(score, 1.0)

def check_vlc_play_and_stop__23bd9b1b573f33533486832ddf9f8705(actual_config_path, rule):
    """
    Checks if VLC's 'Play and stop' playlist setting (play-and-stop) is enabled.
    Values: 0=disabled (default), 1=enabled
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_play_and_stop = rule.get('expected_play_and_stop', '1')
    if isinstance(expected_play_and_stop, int):
        expected_play_and_stop = str(expected_play_and_stop)
    try:
        play_and_stop = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'play-and-stop=' in line:
                play_and_stop = line.split('=')[-1].strip()
        if play_and_stop == expected_play_and_stop:
            return 1.0
        else:
            logger.warning(f'play-and-stop mismatch: expected={expected_play_and_stop}, actual={play_and_stop}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_audio_extraction__e0d33b8eb904792c666a0fe3a33dc4f5(result, expected, **options):
    """Check audio extraction result with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('file_exists'):
        score += 0.34
    if result.get('is_audio'):
        score += 0.33
    if result.get('has_content'):
        score += 0.33
    return min(score, 1.0)

def check_audio_format__3c3196d8d3d24397883e7bbcfddc8dd0(result, expected, **options):
    """
    Check if audio file exists with expected format and properties.
    Partial credit:
      0.3 - File exists
      0.3 - Audio-only (has audio, no video)
      0.2 - Correct audio format/codec
      0.2 - Reasonable duration
    """
    if not result or not isinstance(result, dict) or (not result.get('exists')):
        return 0.0
    if result.get('parse_error'):
        return 0.3
    score = 0.0
    score += 0.3
    if result.get('has_audio') and (not result.get('has_video')):
        score += 0.3
    elif result.get('has_audio'):
        score += 0.15
    expected_formats = expected.get('expected_formats', [])
    expected_codecs = expected.get('expected_codecs', [])
    actual_format = result.get('format_name', '')
    actual_codec = result.get('audio_codec', '')
    format_match = any((f in actual_format for f in expected_formats)) if expected_formats else False
    codec_match = any((c == actual_codec for c in expected_codecs)) if expected_codecs else False
    if format_match or codec_match:
        score += 0.2
    duration = result.get('duration', 0)
    min_duration = expected.get('min_duration', 100)
    max_duration = expected.get('max_duration', 300)
    if min_duration <= duration <= max_duration:
        score += 0.2
    return min(score, 1.0)

def check_vlc_snapshot_file__da77e4abbc92d5a7a2f84a9516d6d0b6(result, expected, **options):
    """Check if snapshot file exists at expected path and is a valid image."""
    if not isinstance(result, dict) or not result.get('exists', False):
        return 0.0
    score = 0.0
    score += 0.5
    mime_type = result.get('mime_type', '')
    if mime_type.startswith('image/'):
        score += 0.3
    file_size = result.get('file_size', 0)
    if file_size > 1024:
        score += 0.2
    return min(score, 1.0)

def check_audio_format__9e6e6e9addabc1fa910c5d7bd33b2f51(result, expected, **options):
    """
    Check if audio file exists with expected format and properties.
    Partial credit:
      0.3 - File exists
      0.3 - Audio-only (has audio, no video)
      0.2 - Correct audio format/codec
      0.2 - Reasonable duration
    """
    if not result or not isinstance(result, dict) or (not result.get('exists')):
        return 0.0
    if result.get('parse_error'):
        return 0.3
    score = 0.0
    score += 0.3
    if result.get('has_audio') and (not result.get('has_video')):
        score += 0.3
    elif result.get('has_audio'):
        score += 0.15
    expected_formats = expected.get('expected_formats', [])
    expected_codecs = expected.get('expected_codecs', [])
    actual_format = result.get('format_name', '')
    actual_codec = result.get('audio_codec', '')
    format_match = any((f in actual_format for f in expected_formats)) if expected_formats else False
    codec_match = any((c == actual_codec for c in expected_codecs)) if expected_codecs else False
    if format_match or codec_match:
        score += 0.2
    duration = result.get('duration', 0)
    min_duration = expected.get('min_duration', 100)
    max_duration = expected.get('max_duration', 300)
    if min_duration <= duration <= max_duration:
        score += 0.2
    return min(score, 1.0)

def check_video_dimensions__92a1bbba6e28a320b340946ad5c7a75d(result, expected, **options):
    """Check video dimensions match expected values with partial credit."""
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    score += 0.5
    exp_w = expected.get('expected_width')
    exp_h = expected.get('expected_height')
    if exp_w is not None and exp_h is not None:
        if result.get('width') == exp_w and result.get('height') == exp_h:
            score += 0.5
    return min(score, 1.0)

def check_vlc_dual_snapshot__1118aad2dd728a496771a4813a51c0c2(result, expected, **options):
    """Check if both snapshot files exist and are valid images. Partial credit per file."""
    if not isinstance(result, dict):
        return 0.0
    total_files = expected.get('total_files', 2)
    score = 0.0
    per_file_score = 1.0 / total_files
    for i in range(total_files):
        key = f'file_{i}'
        file_info = result.get(key, {})
        if not isinstance(file_info, dict):
            continue
        if file_info.get('exists', False):
            file_score = 0.0
            file_score += 0.6
            mime_type = file_info.get('mime_type', '')
            if mime_type.startswith('image/'):
                file_score += 0.2
            file_size = file_info.get('file_size', 0)
            if file_size > 1024:
                file_score += 0.2
            score += per_file_score * file_score
    return min(score, 1.0)

def check_audio_clip_format__7c8156c2a09b3bba70e2cfa0aa1a4dcc(result, expected, **options):
    """
    Check if audio clip file exists with expected format and approximate duration.
    Partial credit:
      0.3 - File exists
      0.3 - Audio-only (has audio, no video)
      0.2 - Correct audio format/codec
      0.2 - Duration within expected range (clip length)
    """
    if not result or not isinstance(result, dict) or (not result.get('exists')):
        return 0.0
    if result.get('parse_error'):
        return 0.3
    score = 0.0
    score += 0.3
    if result.get('has_audio') and (not result.get('has_video')):
        score += 0.3
    elif result.get('has_audio'):
        score += 0.15
    expected_formats = expected.get('expected_formats', [])
    expected_codecs = expected.get('expected_codecs', [])
    actual_format = result.get('format_name', '')
    actual_codec = result.get('audio_codec', '')
    format_match = any((f in actual_format for f in expected_formats)) if expected_formats else False
    codec_match = any((c == actual_codec for c in expected_codecs)) if expected_codecs else False
    if format_match or codec_match:
        score += 0.2
    duration = result.get('duration', 0)
    min_duration = expected.get('min_duration', 50)
    max_duration = expected.get('max_duration', 70)
    if min_duration <= duration <= max_duration:
        score += 0.2
    return min(score, 1.0)

def check_video_dimensions__22c3b2acb75b2231bd0ecc7a3dd0e178(result, expected, **options):
    """Check video dimensions match expected values with partial credit."""
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    score += 0.5
    exp_w = expected.get('expected_width')
    exp_h = expected.get('expected_height')
    if exp_w is not None and exp_h is not None:
        if result.get('width') == exp_w and result.get('height') == exp_h:
            score += 0.5
    return min(score, 1.0)

def check_wav_file_exists__1606287de635101a5a1b8cd46a5bf405(result, expected, **options):
    """Check if a WAV file exists and has a valid size.

    Args:
        result: output from vm_command_line getter (file size string or NOT_FOUND)
        expected: rules dict with 'min_size' (minimum file size in bytes)
    Returns:
        float: 1.0 if file exists and meets min size, 0.0 otherwise
    """
    if result is None or 'NOT_FOUND' in str(result):
        return 0.0
    try:
        file_size = int(str(result).strip())
        min_size = int(expected.get('min_size', 1000))
        if file_size >= min_size:
            return 1.0
        return 0.0
    except (ValueError, TypeError):
        return 0.0

def check_vlc_snapshot_exists__786ac827ab684b501ddc668a08ec2cfa(result, expected, **options):
    """Check if a VLC snapshot file was created."""
    if isinstance(result, dict) and result.get('snapshot_exists', False):
        return 1.0
    return 0.0

def check_vlc_snapshot_exists__441247156205cdab6fb75baec6278a39(result, expected, **options):
    """Check if VLC snapshot was taken successfully."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False) and result.get('count', 0) >= expected.get('min_count', 1):
        return 1.0
    return 0.0

def check_vlc_network_cache__72b69002c9e414b44114ece1b158222e(actual_config_path, rule, **options):
    """Check if VLC's network-caching is set to the expected value."""
    try:
        with open(actual_config_path, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception:
        return 0.0
    expected_value = str(rule.get('expected_cache_ms', ''))
    for line in config_file.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'network-caching' in line:
            current_value = line.split('=', 1)[-1].strip()
            if current_value == expected_value:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_vlc_qt_continue__e543a11713a93c5fff74c1d7c5501303(actual_config_path, rule):
    """
    Checks if VLC's 'Continue playback?' setting (qt-continue) matches the expected value.
    Values: 0=Never, 1=Ask (default), 2=Always
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_continue = rule.get('expected_qt_continue', '2')
    if isinstance(expected_qt_continue, int):
        expected_qt_continue = str(expected_qt_continue)
    try:
        qt_continue = '1'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'qt-continue=' in line:
                qt_continue = line.split('=')[-1].strip()
        if qt_continue == expected_qt_continue:
            return 1.0
        else:
            logger.warning(f'qt-continue mismatch: expected={expected_qt_continue}, actual={qt_continue}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_duration__38add8a045487bc69861547adf883b78_qw35sft2_3c8ca599(result, expected, **options):
    """Return 1.0 if video exists and duration falls within [min_duration, max_duration]."""
    if not isinstance(result, dict) or not result.get('exists', False):
        return 0.0
    duration = result.get('duration', -1.0)
    if duration < 0:
        return 0.0
    min_dur = expected.get('min_duration', 2.5)
    max_dur = expected.get('max_duration', 3.5)
    return 1.0 if min_dur <= duration <= max_dur else 0.0

def check_video_duration__dd6b82d69d19c741223e8dfd4c176567_qw35sft2_d3a22df1(result, expected, **options):
    """Return 1.0 if video exists and duration falls within [min_duration, max_duration]."""
    if not isinstance(result, dict) or not result.get('exists', False):
        return 0.0
    duration = result.get('duration', -1.0)
    if duration < 0:
        return 0.0
    min_dur = expected.get('min_duration', 1.5)
    max_dur = expected.get('max_duration', 2.5)
    return 1.0 if min_dur <= duration <= max_dur else 0.0

def check_single_mp3_tags__85f79002782d5324b058fe59de70692f_qw35sft2_8c651e51(result, expected, **options):
    """
    Check that a single MP3 file has the expected title and artist tags.

    result: stdout string from vm_command_line running a mutagen probe script.
            Expected format (lines):
                title=<value>
                artist=<value>
    expected: dict already unwrapped from 'rules', e.g.
              {"title": "Nights in Shanghai", "artist": "Zhou Xuan"}
    Returns 0.5 per correct tag (total 1.0 if both match).
    """
    if not isinstance(result, str) or not result.strip():
        return 0.0
    tags = {}
    for line in result.strip().splitlines():
        if '=' in line:
            key, _, value = line.partition('=')
            tags[key.strip().lower()] = value.strip()
    exp_title = expected.get('title', '').strip().lower()
    exp_artist = expected.get('artist', '').strip().lower()
    actual_title = tags.get('title', '').lower()
    actual_artist = tags.get('artist', '').lower()
    score = 0.0
    if exp_title and actual_title == exp_title:
        score += 0.5
    if exp_artist and actual_artist == exp_artist:
        score += 0.5
    return score

def check_single_mp3_tags__367848f954ef0bdc8124c00e1a0efa09_qw35sft2_7fc7746c(result, expected, **options):
    """
    Check that a single MP3 file has the expected title and artist tags.

    result: stdout string from vm_command_line running a mutagen probe script.
            Expected format (lines):
                title=<value>
                artist=<value>
    expected: dict already unwrapped from 'rules', e.g.
              {"title": "Tears of Dancing Girl", "artist": "Han Baoyi"}
    Returns 0.5 per correct tag (total 1.0 if both match).
    """
    if not isinstance(result, str) or not result.strip():
        return 0.0
    tags = {}
    for line in result.strip().splitlines():
        if '=' in line:
            key, _, value = line.partition('=')
            tags[key.strip().lower()] = value.strip()
    exp_title = expected.get('title', '').strip().lower()
    exp_artist = expected.get('artist', '').strip().lower()
    actual_title = tags.get('title', '').lower()
    actual_artist = tags.get('artist', '').lower()
    score = 0.0
    if exp_title and actual_title == exp_title:
        score += 0.5
    if exp_artist and actual_artist == exp_artist:
        score += 0.5
    return score

def check_vlc_fullscreen_config__d319a77e7f319beb729da2938a21127c_qw35sft2_875bb6eb(result, expected, **options):
    """
    Check if VLC fullscreen setting is configured as expected by parsing the vlcrc file.

    Args:
        result: Local file path string to the vlcrc config file (from vlc_config getter).
        expected: Rules dict, e.g. {'fullscreen': True} to require fullscreen=1,
                  or {'fullscreen': False} to require fullscreen=0.
    Returns:
        1.0 if the fullscreen setting matches expected, 0.0 otherwise.
    """
    if not result or not isinstance(result, str):
        return 0.0
    want_fullscreen = expected.get('fullscreen', True)
    expected_value = '1' if want_fullscreen else '0'
    try:
        with open(result, 'r', encoding='utf-8', errors='replace') as f:
            for line in f:
                stripped = line.strip()
                if stripped.startswith('#'):
                    continue
                if stripped.startswith('fullscreen='):
                    actual_value = stripped.split('=', 1)[1].strip()
                    return 1.0 if actual_value == expected_value else 0.0
    except Exception:
        return 0.0
    return 0.0

def check_single_mp3_tags__d208efee5ee527e892400da003674774_qw35sft2_41c19618(result, expected, **options):
    """
    Check that a single MP3 file has the expected title and artist tags.

    result: stdout string from vm_command_line running a mutagen probe script.
            Expected format (lines):
                title=<value>
                artist=<value>
    expected: dict already unwrapped from 'rules', e.g.
              {"title": "Red Daughter", "artist": "Chen Shaohua"}
    Returns 0.5 per correct tag (total 1.0 if both match).
    """
    if not isinstance(result, str) or not result.strip():
        return 0.0
    tags = {}
    for line in result.strip().splitlines():
        if '=' in line:
            key, _, value = line.partition('=')
            tags[key.strip().lower()] = value.strip()
    exp_title = expected.get('title', '').strip().lower()
    exp_artist = expected.get('artist', '').strip().lower()
    actual_title = tags.get('title', '').lower()
    actual_artist = tags.get('artist', '').lower()
    score = 0.0
    if exp_title and actual_title == exp_title:
        score += 0.5
    if exp_artist and actual_artist == exp_artist:
        score += 0.5
    return score

def check_mp3_exists_mp4_removed__e17edfd562e8913a3a72a3bcb87c9c97_qw35sft2_095d9ab3(result, expected, **options):
    """Return partial credit: 0.5 for MP3 created, 0.5 for original MP4 deleted.
    Initial state (no MP3, MP4 exists) yields 0.0 (negative-control resistant).
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('mp3_exists'):
        score += 0.5
    if not result.get('mp4_exists'):
        score += 0.5
    return score

def check_vlc_record_path__33b97504f4caa5a543c607fe8892e403_qw35sft2_291ff20a(result, expected, **options):
    """Check if VLC recording path matches expected path."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_path = (result.get('record_path') or '').rstrip('/')
    expected_path = (expected.get('expected_record_path') or '').rstrip('/')
    if not expected_path:
        return 0.0
    return 1.0 if actual_path == expected_path else 0.0

def check_vlc_max_volume__95b68fd0ad40f238d037990e09029619_qw35sft2_a93a4cf8(result, expected, **options):
    """Check that VLC qt-max-volume matches the expected value.

    expected (already unwrapped from rules by get_rule()):
        expected_max_volume: int or str, e.g. 200
    """
    if not isinstance(result, dict):
        return 0.0
    expected_val = str(expected.get('expected_max_volume', ''))
    actual_val = str(result.get('qt_max_volume', ''))
    if actual_val == expected_val:
        logger_qw35sft2_a1dd18.info('VLC max volume check: PASS (qt-max-volume=%s)', actual_val)
        return 1.0
    logger_qw35sft2_a1dd18.info('VLC max volume check: FAIL (expected=%s, got=%s)', expected_val, actual_val)
    return 0.0

def check_vlc_dual_config__de221324409107598f868e822754ff7f_qw35sft2_fa346a76(result, expected, **options):
    """
    Checks two VLC config objectives with partial credit (0.5 each):
      1. qt-minimal-view == expected_qt_minimal_view  (0.5 pts)
      2. qt-max-volume   == expected_qt_max_volume    (0.5 pts)

    expected is already unwrapped by get_rule(), so access keys directly.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_minimal = str(expected.get('expected_qt_minimal_view', '1'))
    actual_minimal = result.get('qt_minimal_view', '0')
    if actual_minimal == expected_minimal:
        score += 0.5
    expected_volume = str(expected.get('expected_qt_max_volume', '200'))
    actual_volume = result.get('qt_max_volume', '125')
    if actual_volume == expected_volume:
        score += 0.5
    return score

def check_vlc_snap_and_rename__25f0d8c371fba607f6a6c0535521c207_qw35sft2_f56f74de(result, expected, **options):
    """Partial credit: 0.5 for snapshot on Desktop, 0.5 for video renamed to trailer.mp4."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('snapshot_on_desktop'):
        score += 0.5
    if result.get('trailer_renamed') and result.get('original_removed'):
        score += 0.5
    elif result.get('trailer_renamed'):
        score += 0.25
    return min(score, 1.0)

def check_vlc_play_and_random__cd6a6657f06ad537781a5ce037c2b82a_qw35sft2_ab28cbbc(result, expected, **options):
    """Check VLC is playing the correct file AND random/shuffle mode is enabled.

    Partial credit: 0.5 for playing correct file, 0.5 for random enabled.
    expected (already unwrapped from rules by get_rule()):
        expected_file_name: str
        expected_random: str, e.g. "true"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_file = str(expected.get('expected_file_name', ''))
    actual_state = str(result.get('state', ''))
    actual_file = str(result.get('file_name', '') or '')
    if actual_state == 'playing' and expected_file and (expected_file in actual_file):
        score += 0.5
    expected_random = str(expected.get('expected_random', 'true')).lower()
    actual_random = str(result.get('random', 'false')).lower()
    if actual_random == expected_random:
        score += 0.5
    return min(score, 1.0)

def check_vlc_brightness__1e9b265786c948c671445fe3130d0b36_qw35sft2_1069eca8(result, expected, **options):
    """Return partial credit score: 0.5 for adjust enabled + 0.5 for brightness increased."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('adjust_enabled'):
        score += 0.5
    brightness = result.get('brightness', _DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f)
    if brightness > _DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f + _BRIGHTNESS_EPSILON_qw35sft2_fa3a1f:
        score += 0.5
    return score

def check_vlc_stream_and_cache__4f78681b35e8e22495902169ef7b5c87_qw35sft2_680c81a3(result, expected, **options):
    """Check that the Apple HLS stream was opened (0.5) and network caching is set correctly (0.5).

    expected (already unwrapped from rules):
        expected_url: str - URL that should appear in VLC's recent MRL list
        expected_network_caching: str - expected network-caching value in ms, e.g. "3000"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_url = expected.get('expected_url', '')
    recent_mrl = result.get('recent_mrl', '')
    if expected_url and expected_url in recent_mrl:
        score += 0.5
    expected_cache = str(expected.get('expected_network_caching', '3000'))
    actual_cache = str(result.get('network-caching', '1000'))
    if actual_cache == expected_cache:
        score += 0.5
    return min(score, 1.0)

def check_snapshot_and_wallpaper__1e243c15dcedcf4e9a4cacc616a0598b_qw35sft2_a4f44df9(result, expected, **options):
    """Partial credit: 0.5 for VLC snapshot in Pictures, 0.5 for wallpaper set to vlcsnap."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    snap_count = result.get('snapshot_count', 0)
    if snap_count >= expected.get('min_snapshot_count', 1):
        score += 0.5
    uri = result.get('wallpaper_uri', '')
    pattern = expected.get('wallpaper_pattern', 'vlcsnap')
    if pattern in uri and ('file://' in uri or uri.startswith('/')):
        score += 0.5
    return min(score, 1.0)

def check_vlc_recfolder_and_maxvol__517e4427a675aa3407c571ed7524794f_qw35sft2_d7155b18(result, expected, **options):
    """
    Partial credit: 0.5 for correct recording folder, 0.5 for correct max volume in vlcrc.
    expected is the already-unwrapped rules dict.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_path = expected.get('expected_record_path', '/home/user/Desktop')
    actual_path = result.get('input_record_path')
    if actual_path == expected_path:
        score += 0.5
        logger_qw35sft2_462da1.info('VLC recording folder check: PASS (path=%s) (+0.5)', actual_path)
    else:
        logger_qw35sft2_462da1.info('VLC recording folder check: FAIL (expected=%s, got=%s)', expected_path, actual_path)
    expected_vol = str(expected.get('expected_max_volume', '150'))
    actual_vol = str(result.get('qt_max_volume', ''))
    if actual_vol == expected_vol:
        score += 0.5
        logger_qw35sft2_462da1.info('VLC max-volume check: PASS (volume=%s) (+0.5)', actual_vol)
    else:
        logger_qw35sft2_462da1.info('VLC max-volume check: FAIL (expected=%s, got=%s)', expected_vol, actual_vol)
    return score

def check_vlc_record_path__e488a41c13481ad2b017a00cb14067ec_qw35sft2_8dd20fab(result, expected, **options):
    """Check if VLC recording path matches expected path."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_path = (result.get('record_path') or '').rstrip('/')
    expected_path = (expected.get('expected_record_path') or '').rstrip('/')
    if not expected_path:
        return 0.0
    return 1.0 if actual_path == expected_path else 0.0

def check_vlc_max_volume__15cb17bd2e24428fe6d640ea8d5d556a_qw35sft2_af4230f7(result, expected, **options):
    """Check that VLC qt-max-volume matches the expected value.

    expected (already unwrapped from rules by get_rule()):
        expected_max_volume: int or str, e.g. 300
    """
    if not isinstance(result, dict):
        return 0.0
    expected_val = str(expected.get('expected_max_volume', ''))
    actual_val = str(result.get('qt_max_volume', ''))
    if actual_val == expected_val:
        logger_qw35sft2_a878d7.info('VLC max volume check: PASS (qt-max-volume=%s)', actual_val)
        return 1.0
    logger_qw35sft2_a878d7.info('VLC max volume check: FAIL (expected=%s, got=%s)', expected_val, actual_val)
    return 0.0

def check_vlc_instance_settings__71c2e6dc1ced259311cab9f07521fd96_qw35sft2_65756d41(result, expected, **options):
    """Check that both VLC single-instance settings are disabled (=0).

    Scoring:
      +0.5  'Allow only one instance' (one-instance) == 0
      +0.5  'Use only one instance when started from file manager' (one-instance-when-started-from-file) == 0
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if str(result.get('one_instance', 'not_set')) == '0':
        score += 0.5
    if str(result.get('one_instance_from_file', 'not_set')) == '0':
        score += 0.5
    return score

def check_mp3_validity__35ef4e9e6e88fc273aff9d3f203be720_qw35sft2_54cf8bf4(result, expected, **options):
    """Return 1.0 if the file has a valid MP3 magic header and is at least min_size bytes."""
    if not isinstance(result, dict):
        return 0.0
    min_size = expected.get('min_size', 100000)
    if result.get('valid') and result.get('size', 0) >= min_size:
        return 1.0
    return 0.0

def check_vlc_play_and_rate__a87ab91c6345361c3c6e374fc52958bf_qw35sft2_949dbd4a(result, expected, **options):
    """Check VLC is playing the correct file AND playback rate matches expected.

    Partial credit: 0.5 for playing correct file, 0.5 for correct rate.
    expected (already unwrapped from rules by get_rule()):
        expected_file_name: str
        expected_rate: float, e.g. 2.0
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_file = str(expected.get('expected_file_name', ''))
    actual_state = str(result.get('state', ''))
    actual_file = str(result.get('file_name', '') or '')
    if actual_state == 'playing' and expected_file and (expected_file in actual_file):
        score += 0.5
    expected_rate = float(expected.get('expected_rate', 2.0))
    try:
        actual_rate = float(result.get('rate', '1'))
        if abs(actual_rate - expected_rate) < 0.1:
            score += 0.5
    except (ValueError, TypeError):
        pass
    return min(score, 1.0)

def check_vlc_rename__bd8f719f06ca22a9a2d291dd693db756_qw35sft2_53a730de(result, expected, **options):
    """
    Partial-credit check for a file rename on the Desktop.

    result  – string from get_vlc_rename_check__bd8f719f06ca22a9a2d291dd693db756,
              e.g. 'new_exists=1 old_removed=1'
    expected – rules dict (already unwrapped by get_rule).
               Keys used: 'new_exists' (default '1'), 'old_removed' (default '1').

    Scoring:
      +0.5  if new file '1984_Apple_Macintosh_Commercial.mp4' exists
      +0.5  if old file 'flipped_1984_Apple_Macintosh_Commercial.mp4' is absent
    """
    if not result or result == 'error':
        logger_qw35sft2_211cdf.warning("check_vlc_rename: bad result '%s'", result)
        return 0.0
    score = 0.0
    if 'new_exists=1' in result:
        score += 0.5
    if 'old_removed=1' in result:
        score += 0.5
    logger_qw35sft2_211cdf.info('check_vlc_rename: result=%r score=%.1f', result, score)
    return score

def check_vlc_cone_and_oneinstance__6535f71afc4db044467cc2ad41010df9_qw35sft2_582f829f(result, expected, **options):
    """Check that background cone is disabled AND one-instance-when-started-from-file is disabled.

    Partial credit: 0.5 per sub-goal.
    expected (already unwrapped from rules):
        expected_qt_bgcone: str, e.g. "0"
        expected_one_instance: str, e.g. "0"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_cone = str(expected.get('expected_qt_bgcone', '0'))
    actual_cone = str(result.get('qt-bgcone', '1'))
    if actual_cone == expected_cone:
        score += 0.5
    expected_oi = str(expected.get('expected_one_instance', '0'))
    actual_oi = str(result.get('one-instance-when-started-from-file', '1'))
    if actual_oi == expected_oi:
        score += 0.5
    return min(score, 1.0)

def check_wallpaper_and_desktop_snapshot__5688825554c17cb3fb1e24246b369618_qw35sft2_53323222(result, expected, **options):
    """Partial credit: 0.5 for wallpaper set to vlcsnap, 0.5 for snapshot on Desktop."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    uri = result.get('wallpaper_uri', '')
    pattern = expected.get('wallpaper_pattern', 'vlcsnap')
    if pattern in uri and ('file://' in uri or uri.startswith('/')):
        score += 0.5
    desk_count = result.get('desktop_snapshot_count', 0)
    if desk_count >= expected.get('min_desktop_count', 1):
        score += 0.5
    return min(score, 1.0)

def check_vlc_stream_and_cone__e0203e1b79d3816d60eca3abd8c45b33_qw35sft2_6d2cf461(result, expected, **options):
    """Check that the Apple HLS stream was opened (0.5) and background cone is disabled (0.5).

    expected (already unwrapped from rules):
        expected_url: str - URL that should appear in VLC's recent MRL list
        expected_qt_bgcone: str - expected qt-bgcone value, e.g. "0"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_url = expected.get('expected_url', '')
    recent_mrl = result.get('recent_mrl', '')
    if expected_url and expected_url in recent_mrl:
        score += 0.5
    expected_cone = str(expected.get('expected_qt_bgcone', '0'))
    actual_cone = str(result.get('qt-bgcone', '1'))
    if actual_cone == expected_cone:
        score += 0.5
    return min(score, 1.0)

def check_vlc_snap_in_captures__6d219cf2312d28ed73dcf8a9fc508b4d_qw35sft2_288b20f7(result, expected, **options):
    """Partial credit: 0.3 for captures/ folder, 0.7 for interstellar.png inside it."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('captures_folder_exists'):
        score += 0.3
    if result.get('interstellar_in_captures'):
        score += 0.7
    return score

def check_vlc_fullscreen_and_snapshot__b80790f2f04245c0798c2ba24e3ceaf2_qw35sft2_5c114c55(result, expected, **options):
    """
    Partial credit: 0.5 for VLC fullscreen, 0.5 for at least one snapshot PNG in ~/Pictures.
    expected is the already-unwrapped rules dict.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('is_fullscreen'):
        score += 0.5
        logger_qw35sft2_ccdace.info('VLC fullscreen check: PASS (+0.5)')
    else:
        logger_qw35sft2_ccdace.info('VLC fullscreen check: FAIL')
    if result.get('snapshot_count', 0) > 0:
        score += 0.5
        logger_qw35sft2_ccdace.info('VLC snapshot check: PASS (%d PNG files found) (+0.5)', result.get('snapshot_count'))
    else:
        logger_qw35sft2_ccdace.info('VLC snapshot check: FAIL (no PNG files found)')
    return score

def check_vlc_maxvol_and_bgcone__f99d034b7e5202fd562c023da2142c74_qw35sft2_66853c69(result, expected, **options):
    """Check max volume AND background cone settings with partial credit.

    Partial credit: 0.5 for correct max volume, 0.5 for bgcone disabled (0).
    expected (already unwrapped from rules by get_rule()):
        expected_max_volume: int or str, e.g. 200
        expected_bgcone: str, e.g. "0" (disabled)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_vol = str(expected.get('expected_max_volume', ''))
    actual_vol = str(result.get('qt_max_volume', ''))
    if actual_vol == expected_vol:
        score += 0.5
        logger_qw35sft2_544f1c.info('VLC max volume check: PASS (qt-max-volume=%s) (+0.5)', actual_vol)
    else:
        logger_qw35sft2_544f1c.info('VLC max volume check: FAIL (expected=%s, got=%s)', expected_vol, actual_vol)
    expected_bgcone = str(expected.get('expected_bgcone', '0'))
    actual_bgcone = str(result.get('qt_bgcone', ''))
    if actual_bgcone == expected_bgcone:
        score += 0.5
        logger_qw35sft2_544f1c.info('VLC bgcone check: PASS (qt-bgcone=%s) (+0.5)', actual_bgcone)
    else:
        logger_qw35sft2_544f1c.info('VLC bgcone check: FAIL (expected=%s, got=%s)', expected_bgcone, actual_bgcone)
    return score

def check_vlc_playback_prefs__db101724f343bdcc09eb9f11a0e77c94_qw35sft2_55f6826f(result, expected, **options):
    """Check VLC file-manager instance setting and continue-playback setting.

    Scoring:
      +0.5  'Use only one instance when started from file manager' == 0 (disabled)
      +0.5  'Continue playback?' == 2 (Always resume without asking)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if str(result.get('one_instance_from_file', 'not_set')) == '0':
        score += 0.5
    if str(result.get('continue_playback', 'not_set')) == '2':
        score += 0.5
    return score

def check_vlc_effects_dialog__80d9e76aa584c867ce735bfda2f705a1_qw35sft2_b960a33e(result, expected, **options):
    """
    Check whether VLC's 'Adjustments and Effects' dialog is currently open
    by searching the accessibility tree string for the dialog title tokens.

    result   – accessibility-tree string from get_vlc_effects_dialog__80d9e76...
    expected – rules dict (already unwrapped by get_rule).
               Optional key 'dialog_tokens' (list[str]) to search for;
               defaults to ['Adjustments', 'Effects'].

    Returns 1.0 if all tokens are found in the tree, 0.0 otherwise.
    """
    if not result:
        logger_qw35sft2_1be130.warning('check_vlc_effects_dialog: empty accessibility tree')
        return 0.0
    tokens = expected.get('dialog_tokens', ['Adjustments', 'Effects'])
    if all((tok in result for tok in tokens)):
        logger_qw35sft2_1be130.info('check_vlc_effects_dialog: dialog found (tokens=%s)', tokens)
        return 1.0
    logger_qw35sft2_1be130.info('check_vlc_effects_dialog: dialog NOT found. tokens=%s, tree_snippet=%.200s', tokens, result)
    return 0.0

def check_vlc_play_and_repeat__dc3aa9f060ebd375fbf717e1b7e05e93_qw35sft2_b1cfa8c7(result, expected, **options):
    """Check VLC is playing the correct file AND repeat mode is enabled.

    Partial credit: 0.5 for playing correct file, 0.5 for repeat enabled.
    expected (already unwrapped from rules by get_rule()):
        expected_file_name: str
        expected_repeat: str, e.g. "true"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_file = str(expected.get('expected_file_name', ''))
    actual_state = str(result.get('state', ''))
    actual_file = str(result.get('file_name', '') or '')
    if actual_state == 'playing' and expected_file and (expected_file in actual_file):
        score += 0.5
    expected_repeat = str(expected.get('expected_repeat', 'true')).lower()
    actual_repeat = str(result.get('repeat', 'false')).lower()
    if actual_repeat == expected_repeat:
        score += 0.5
    return min(score, 1.0)

def check_mp3_mp4_both_exist__41a079bbfac3c55bed03337726ec0862_qw35sft2_251d7a95(result, expected, **options):
    """Non-destructive conversion check.
    - 1.0 if MP3 was created AND MP4 still exists (perfect non-destructive conversion)
    - 0.5 if MP3 exists but MP4 was deleted (conversion done but original removed)
    - 0.0 if MP3 does not exist (conversion not completed)
    Initial state (no MP3, MP4 exists) yields 0.0 (negative-control resistant).
    """
    if not isinstance(result, dict):
        return 0.0
    mp3_exists = result.get('mp3_exists', False)
    mp4_exists = result.get('mp4_exists', False)
    if not mp3_exists:
        return 0.0
    if mp4_exists:
        return 1.0
    return 0.5

def check_vlc_prefs__8b6f202a2cac4c2b2f141e2f0645d358_qw35sft2_e766877d(result, expected, **options):
    """Check VLC recording path (0.5) and video-title-show disabled (0.5)."""
    if isinstance(result, dict) and result.get('error') and (not result.get('record_path')) and (not result.get('video_title_show')):
        return 0.0
    score = 0.0
    actual_path = (result.get('record_path') or '').rstrip('/')
    expected_path = (expected.get('expected_record_path') or '').rstrip('/')
    if expected_path and actual_path == expected_path:
        score += 0.5
    video_title_show = result.get('video_title_show')
    if video_title_show is not None and str(video_title_show).strip() == '0':
        score += 0.5
    return min(score, 1.0)

def check_vlc_saturation__493da9b8d1a20a26a98a5cc60b751e43_qw35sft2_fdfb9d9d(result, expected, **options):
    """Return partial credit: 0.5 for adjust enabled + 0.5 for saturation set to minimum."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('adjust_enabled'):
        score += 0.5
    saturation = result.get('saturation', 1.0)
    if saturation <= _SATURATION_GRAYSCALE_MAX_qw35sft2_7934da:
        score += 0.5
    return score

def check_snapshot_file_exists__4e0abb3a76e2d3729689698c9288f79c_qw35sft2_402eb8d2(result, expected, **options):
    """Check that at least one VLC snapshot PNG was saved in Pictures."""
    if not isinstance(result, dict):
        return 0.0
    count = result.get('snapshot_count', 0)
    min_count = expected.get('min_count', 1)
    return 1.0 if count >= min_count else 0.0

def check_vlc_cone_and_volume__78f02ce8f924783e3007881645e6184a_qw35sft2_8ddca2e6(result, expected, **options):
    """Check that background cone is disabled AND max volume is set to expected value.

    Partial credit: 0.5 per sub-goal.
    expected (already unwrapped from rules):
        expected_qt_bgcone: str, e.g. "0"
        expected_qt_max_volume: str, e.g. "200"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_cone = str(expected.get('expected_qt_bgcone', '0'))
    actual_cone = str(result.get('qt-bgcone', '1'))
    if actual_cone == expected_cone:
        score += 0.5
    expected_vol = str(expected.get('expected_qt_max_volume', '200'))
    actual_vol = str(result.get('qt-max-volume', '125'))
    if actual_vol == expected_vol:
        score += 0.5
    return min(score, 1.0)

def check_vlc_stream_and_snapshot__b5aec674ad8f49621ab70bae692967c5_qw35sft2_331ca700(result, expected, **options):
    """Check that the Apple HLS stream was opened (0.5) and a snapshot was taken (0.5).

    expected (already unwrapped from rules):
        expected_url: str - URL that should appear in VLC's recent MRL list
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_url = expected.get('expected_url', '')
    recent_mrl = result.get('recent_mrl', '')
    if expected_url and expected_url in recent_mrl:
        score += 0.5
    snapshot_count = result.get('snapshot_count', 0)
    if snapshot_count > 0:
        score += 0.5
    return min(score, 1.0)

def vlc_snapshot_dual_loc__b8a50137decabc772595757ced9c456a_qw35sft2_d6567a65(result, expected, **options):
    """Partial credit: 0.5 for Desktop copy, 0.5 for Pictures copy."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('desktop_has_interstellar'):
        score += 0.5
    if result.get('pictures_has_interstellar'):
        score += 0.5
    return score

def check_vlc_fullscreen_and_maxvol__ecd5fecf55729059a58a10d83c8c9eeb_qw35sft2_de69a72b(result, expected, **options):
    """
    Partial credit: 0.5 for VLC fullscreen, 0.5 for qt-max-volume set to expected value.
    expected is the already-unwrapped rules dict.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('is_fullscreen'):
        score += 0.5
        logger_qw35sft2_a640c9.info('VLC fullscreen check: PASS (+0.5)')
    else:
        logger_qw35sft2_a640c9.info('VLC fullscreen check: FAIL')
    expected_vol = str(expected.get('expected_max_volume', '150'))
    actual_vol = str(result.get('qt_max_volume', ''))
    if actual_vol == expected_vol:
        score += 0.5
        logger_qw35sft2_a640c9.info('VLC max-volume check: PASS (volume=%s) (+0.5)', actual_vol)
    else:
        logger_qw35sft2_a640c9.info('VLC max-volume check: FAIL (expected=%s, got=%s)', expected_vol, actual_vol)
    return score

def check_vlc_triple_prefs__29cde5ba2a87fed510792760a2c64d1f_qw35sft2_b89ee98a(result, expected, **options):
    """Check all three VLC multi-instance and playback settings.

    Scoring:
      +0.34  'Allow only one instance' (one-instance) == 0 (disabled)
      +0.33  'Use only one instance when started from file manager' == 0 (disabled)
      +0.33  'Continue playback?' == 2 (Always resume without asking)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if str(result.get('one_instance', 'not_set')) == '0':
        score += 0.34
    if str(result.get('one_instance_from_file', 'not_set')) == '0':
        score += 0.33
    if str(result.get('continue_playback', 'not_set')) == '2':
        score += 0.33
    return min(score, 1.0)

def check_vlc_play_and_record_path__d9c5d9aac51555091a28197b6da34259_qw35sft2_7910036a(result, expected, **options):
    """Check VLC is playing the correct file AND recording directory is set to expected path.

    Partial credit: 0.5 for playing correct file, 0.5 for correct record path.
    expected (already unwrapped from rules by get_rule()):
        expected_file_name: str
        expected_record_path: str, e.g. "/home/user/Desktop"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_file = str(expected.get('expected_file_name', ''))
    actual_state = str(result.get('state', ''))
    actual_file = str(result.get('file_name', '') or '')
    if actual_state == 'playing' and expected_file and (expected_file in actual_file):
        score += 0.5
    expected_record = str(expected.get('expected_record_path', ''))
    actual_record = str(result.get('record_path', '') or '')
    if expected_record and actual_record == expected_record:
        score += 0.5
    return min(score, 1.0)

def check_mp3_file_size__55ceccfa8167b4a354049a9775b28527_qw35sft2_098ff36e(result, expected, **options):
    """Return 1.0 if Baby Justin Bieber.mp3 exists and is at least min_size bytes."""
    if not isinstance(result, dict):
        return 0.0
    size = result.get('size', 0)
    min_size = expected.get('min_size', 500000)
    return 1.0 if size >= min_size else 0.0

def check_vlc_maxvol_and_bgcone_expands__d25099f13e7da2513cf33c07be0e7955_qw35sft2_378e65f4(result, expected, **options):
    """Check max volume AND expanding-cone animation settings with partial credit.

    Partial credit: 0.5 for correct max volume, 0.5 for bgcone-expands disabled (0).
    expected (already unwrapped from rules by get_rule()):
        expected_max_volume: int or str, e.g. 200
        expected_bgcone_expands: str, e.g. "0" (disabled)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_vol = str(expected.get('expected_max_volume', ''))
    actual_vol = str(result.get('qt_max_volume', ''))
    if actual_vol == expected_vol:
        score += 0.5
        logger_qw35sft2_208a9d.info('VLC max volume check: PASS (qt-max-volume=%s) (+0.5)', actual_vol)
    else:
        logger_qw35sft2_208a9d.info('VLC max volume check: FAIL (expected=%s, got=%s)', expected_vol, actual_vol)
    expected_expands = str(expected.get('expected_bgcone_expands', '0'))
    actual_expands = str(result.get('qt_bgcone_expands', ''))
    if actual_expands == expected_expands:
        score += 0.5
        logger_qw35sft2_208a9d.info('VLC bgcone-expands check: PASS (qt-bgcone-expands=%s) (+0.5)', actual_expands)
    else:
        logger_qw35sft2_208a9d.info('VLC bgcone-expands check: FAIL (expected=%s, got=%s)', expected_expands, actual_expands)
    return score

def check_vlc_threshold__c451d2a7e7a5b8e510cce0c34da4325d_qw35sft2_19eb83a8(result, expected, **options):
    """Return partial credit: 0.5 for adjust enabled + 0.5 for brightness threshold enabled."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('adjust_enabled'):
        score += 0.5
    if result.get('brightness_threshold'):
        score += 0.5
    return score

def check_vlc_cone_and_minview__24dc6e88ff00c69857d313c4738a8f31_qw35sft2_4532da52(result, expected, **options):
    """Check that background cone is disabled AND minimal view is enabled.

    Partial credit: 0.5 per sub-goal.
    expected (already unwrapped from rules):
        expected_qt_bgcone: str, e.g. "0"
        expected_qt_minimal_view: str, e.g. "1"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_cone = str(expected.get('expected_qt_bgcone', '0'))
    actual_cone = str(result.get('qt-bgcone', '1'))
    if actual_cone == expected_cone:
        score += 0.5
    expected_minview = str(expected.get('expected_qt_minimal_view', '1'))
    actual_minview = str(result.get('qt-minimal-view', '0'))
    if actual_minview == expected_minview:
        score += 0.5
    return min(score, 1.0)

def check_vlc_snapshot_wallpaper__d24bdc521260cf2eb34948f145dea565_qw35sft2_66214d9a(result, expected, **options):
    """Check wallpaper URI contains a VLC snapshot file from Pictures folder."""
    if not isinstance(result, dict) or 'wallpaper_uri' not in result:
        return 0.0
    uri = result.get('wallpaper_uri', '')
    if not uri or uri == '':
        return 0.0
    required_pattern = expected.get('required_pattern', 'vlcsnap')
    if required_pattern in uri and ('file://' in uri or uri.startswith('/')):
        return 1.0
    return 0.0

def check_vlc_record_path__f613292e2de2a0e0145aaa24fdebbc72_qw35sft2_67a189da(result, expected, **options):
    """Check if VLC recording path matches expected path."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_path = (result.get('record_path') or '').rstrip('/')
    expected_path = (expected.get('expected_record_path') or '').rstrip('/')
    if not expected_path:
        return 0.0
    return 1.0 if actual_path == expected_path else 0.0

def check_vlc_stream_and_volume__a51ef1127f9b7bededbd54fd205d555d_qw35sft2_6a01a4f9(result, expected, **options):
    """Check that the Apple HLS stream was opened (0.5) and max volume is set correctly (0.5).

    expected (already unwrapped from rules):
        expected_url: str - URL that should appear in VLC's recent MRL list
        expected_qt_max_volume: str - expected qt-max-volume value, e.g. "200"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_url = expected.get('expected_url', '')
    recent_mrl = result.get('recent_mrl', '')
    if expected_url and expected_url in recent_mrl:
        score += 0.5
    expected_vol = str(expected.get('expected_qt_max_volume', '200'))
    actual_vol = str(result.get('qt-max-volume', '125'))
    if actual_vol == expected_vol:
        score += 0.5
    return min(score, 1.0)

def vlc_snapshot_pictures__cdcbbd90330cd55be87f2b9353b8c43c_qw35sft2_1ff63cb9(result, expected, **options):
    """1.0 if scene.png exists in Pictures and is PNG, 0.5 if exists but not PNG."""
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists'):
        return 0.0
    if result.get('is_png') and result.get('size', 0) > 1000:
        return 1.0
    return 0.5

def check_vlc_fullscreen_and_bgcone__986ee683b9dd7fb42919cda3330ee515_qw35sft2_5e67ef02(result, expected, **options):
    """
    Partial credit: 0.5 for VLC fullscreen, 0.5 for qt-bgcone set to expected value (0 = disabled).
    expected is the already-unwrapped rules dict.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('is_fullscreen'):
        score += 0.5
        logger_qw35sft2_0717dc.info('VLC fullscreen check: PASS (+0.5)')
    else:
        logger_qw35sft2_0717dc.info('VLC fullscreen check: FAIL')
    expected_bgcone = str(expected.get('expected_bgcone', '0'))
    actual_bgcone = str(result.get('qt_bgcone', ''))
    if actual_bgcone == expected_bgcone:
        score += 0.5
        logger_qw35sft2_0717dc.info('VLC bgcone check: PASS (qt-bgcone=%s) (+0.5)', actual_bgcone)
    else:
        logger_qw35sft2_0717dc.info('VLC bgcone check: FAIL (expected=%s, got=%s)', expected_bgcone, actual_bgcone)
    return score

def check_vlc_loop_prefs__a02aa930219f8e3a3317ea91cceaa2ff_qw35sft2_b6c8b6bc(result, expected, **options):
    """Check VLC file-manager instance setting and playlist loop setting.

    Scoring:
      +0.5  'Use only one instance when started from file manager' == 0 (disabled)
      +0.5  'Loop' == 1 (playlist loop enabled)
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if str(result.get('one_instance_from_file', 'not_set')) == '0':
        score += 0.5
    if str(result.get('loop', 'not_set')) == '1':
        score += 0.5
    return score

def check_vlc_play_and_loop__c8546449445cd5e2551a333da05cf053_qw35sft2_2ed6f216(result, expected, **options):
    """Check VLC is playing the correct file AND loop mode is enabled.

    Partial credit: 0.5 for playing correct file, 0.5 for loop enabled.
    expected (already unwrapped from rules by get_rule()):
        expected_file_name: str
        expected_loop: str, e.g. "true"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_file = str(expected.get('expected_file_name', ''))
    actual_state = str(result.get('state', ''))
    actual_file = str(result.get('file_name', '') or '')
    if actual_state == 'playing' and expected_file and (expected_file in actual_file):
        score += 0.5
    expected_loop = str(expected.get('expected_loop', 'true')).lower()
    actual_loop = str(result.get('loop', 'false')).lower()
    if actual_loop == expected_loop:
        score += 0.5
    return min(score, 1.0)

def check_vlc_contrast__ab50e89f0aef704ab57d4c9d2579417f_qw35sft2_eac2f79d(result, expected, **options):
    """Return partial credit score: 0.5 for adjust enabled + 0.5 for contrast increased."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    if result.get('adjust_enabled'):
        score += 0.5
    contrast = result.get('contrast', _DEFAULT_CONTRAST_qw35sft2_bffd0c)
    if contrast > _DEFAULT_CONTRAST_qw35sft2_bffd0c + _CONTRAST_EPSILON_qw35sft2_bffd0c:
        score += 0.5
    return score

def check_vlc_max_volume__a9382a8b8a5c0cd01bfb6f50454b9a65_qw35sft2_e32b8685(result, expected, **options):
    """Check that VLC qt-max-volume matches the expected value.

    expected (already unwrapped from rules by get_rule()):
        expected_max_volume: int or str, e.g. 150
    """
    if not isinstance(result, dict):
        return 0.0
    expected_val = str(expected.get('expected_max_volume', ''))
    actual_val = str(result.get('qt_max_volume', ''))
    if actual_val == expected_val:
        logger_qw35sft2_813e96.info('VLC max volume check: PASS (qt-max-volume=%s)', actual_val)
        return 1.0
    logger_qw35sft2_813e96.info('VLC max volume check: FAIL (expected=%s, got=%s)', expected_val, actual_val)
    return 0.0

def check_vlc_cone_and_recordpath__914a6c12022e434d9eb3c729da97ace4_qw35sft2_5861b639(result, expected, **options):
    """Check that background cone is disabled AND recordings path is set to expected directory.

    Partial credit: 0.5 per sub-goal.
    expected (already unwrapped from rules):
        expected_qt_bgcone: str, e.g. "0"
        expected_record_path: str, e.g. "/home/user/Desktop"
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_cone = str(expected.get('expected_qt_bgcone', '0'))
    actual_cone = str(result.get('qt-bgcone', '1'))
    if actual_cone == expected_cone:
        score += 0.5
    expected_path = expected.get('expected_record_path', '')
    actual_path = result.get('input-record-path', '')
    if expected_path and actual_path:
        if os.path.normpath(actual_path) == os.path.normpath(expected_path):
            score += 0.5
    elif expected_path == '' and actual_path == '':
        score += 0.5
    return min(score, 1.0)

def check_mp3_existence__77b822a22aadd202cc6c547979661a88_qw35sft2_2317004e(result, expected, **options):
    """Return 1.0 if Baby Justin Bieber.mp3 exists on the desktop, else 0.0."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('exists') is True else 0.0

def check_vlc_snap_on_desktop__dc35efec6c78cc3245c9e38e3e5c5d4a_qw35sft2_f8227f2f(result, expected, **options):
    """Partial credit: 0.4 for VLC snapshot-path set to Desktop, 0.6 for vlcsnap file on Desktop."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('snap_path_configured'):
        score += 0.4
    if result.get('snap_on_desktop'):
        score += 0.6
    return score

def check_wallpaper_and_vlc_running__2cdc858f118b65a32b6755071b250267_qw35sft2_74c7dcea(result, expected, **options):
    """Partial credit: 0.5 for wallpaper set to vlcsnap, 0.5 for VLC still running."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    uri = result.get('wallpaper_uri', '')
    pattern = expected.get('wallpaper_pattern', 'vlcsnap')
    if pattern in uri and ('file://' in uri or uri.startswith('/')):
        score += 0.5
    vlc_running = result.get('vlc_running', False)
    require_vlc_open = expected.get('require_vlc_open', True)
    if not require_vlc_open or vlc_running:
        score += 0.5
    return min(score, 1.0)

def check_vlc_fullscreen_and_recfolder__27e3d8a3f9225f7e9bc7da9d0f1debd6_qw35sft2_8442bba5(result, expected, **options):
    """
    Partial credit: 0.5 for VLC fullscreen, 0.5 for correct recording folder in vlcrc.
    expected is the already-unwrapped rules dict.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('is_fullscreen'):
        score += 0.5
        logger_qw35sft2_55647c.info('VLC fullscreen check: PASS (+0.5)')
    else:
        logger_qw35sft2_55647c.info('VLC fullscreen check: FAIL')
    expected_path = expected.get('expected_record_path', '/home/user/Desktop')
    actual_path = result.get('input_record_path')
    if actual_path == expected_path:
        score += 0.5
        logger_qw35sft2_55647c.info('VLC recording folder check: PASS (path=%s) (+0.5)', actual_path)
    else:
        logger_qw35sft2_55647c.info('VLC recording folder check: FAIL (expected=%s, got=%s)', expected_path, actual_path)
    return score

def check_vlc_stream_and_ontop__5d8c2e81cb9645b71051b528d54492c6_qw35sft2_47d573db(result, expected, **options):
    """Check that the Apple HLS stream was opened (0.5) and Always on Top is enabled (0.5).

    expected (already unwrapped from rules):
        expected_url: str - URL that should appear in VLC's recent MRL list
        expected_video_on_top: str - expected video-on-top value, "1" = enabled
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_url = expected.get('expected_url', '')
    recent_mrl = result.get('recent_mrl', '')
    if expected_url and expected_url in recent_mrl:
        score += 0.5
    expected_ontop = str(expected.get('expected_video_on_top', '1'))
    actual_ontop = str(result.get('video-on-top', '0'))
    if actual_ontop == expected_ontop:
        score += 0.5
    return min(score, 1.0)

def check_vlc_prefs__a7d5b4b88d0304ec4fb07fa10130fbc0_qw35sft2_0b6de3eb(result, expected, **options):
    """Check VLC recording path (0.5) and video-title-show disabled (0.5)."""
    if isinstance(result, dict) and result.get('error') and (not result.get('record_path')) and (not result.get('video_title_show')):
        return 0.0
    score = 0.0
    actual_path = (result.get('record_path') or '').rstrip('/')
    expected_path = (expected.get('expected_record_path') or '').rstrip('/')
    if expected_path and actual_path == expected_path:
        score += 0.5
    video_title_show = result.get('video_title_show')
    if video_title_show is not None and str(video_title_show).strip() == '0':
        score += 0.5
    return min(score, 1.0)
