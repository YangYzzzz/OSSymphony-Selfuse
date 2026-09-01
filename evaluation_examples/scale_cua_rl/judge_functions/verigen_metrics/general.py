"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_file_copy_and_clipboard__267da608b7e8dd72021770568c2da84f', 'check_all_text_uppercase__0c637a2f0d91975120dc8e0a76554911', 'check_text_file_keywords__ca2a0f6f5509af90b7b3cf524daf3368', 'check_text_replacement__12dc12177d981fb82b80cf3f9509e0cf', 'check_file_listing__6790a81b7f7c76a950ade3c3fcc6cc45', 'check_text_content__ec3e3cd160378737f8c3074a7c1fe7f0', 'check_outline_view__4b3841fb1241cbb16c8a71b3b365dbfe', 'check_download_speed_file__034287361433ba7df6a1a4debdad660e', 'check_python_syntax_valid__94eb7a158102fbe4b2865eae8fd2e0f7', 'check_version_file__4c4a9bd632c3f7487e991bbd8a3f942a', 'check_line_content__c54a7288908815335d3763c41c9cef77', 'check_file_renamed__9b65526186f9fb77c93ff9d7de5c5fcb', 'check_file_exists__0ae8da7e0d709c03ee846a1d011df875', 'check_line_content__350a72c316de89189f7f590d79389d9e', 'check_file_organize__a65eb6fd760478532b3e118be8d079f9', 'check_title_text__1b7b08a28b40e394d9044a5b56ecfe62', 'check_file_moved__e60640f858dbfcce19e818f08302a608', 'check_line_content__c8cb9d1d5c3b002c8ddbeb93d3580891', 'check_first_line__35d6e56c2b031e6f068f1907dd34bebf', 'check_text_uppercase__0a695088802a9cb1473c98bdb4ec1b4f', 'check_text_alignment__a27e67f69ab078486669e5f21f103d16', 'check_file_exists_nonempty__29f89216bfdb55e864e9190350cbf19a', 'check_json_settings_subset__7efbe3197e1cdd53b30f5426c6d455b2', 'check_text_scaling__04c87fb26c99f62f08868d28d96db622', 'check_direct_json_object', 'check_file_dual_location__454f85dd5a1086d0f2418681bdac8477', 'check_file_content__c80ddc5f2ae016c0c4e6f2749cc5143e', 'check_file_rename__5128f87c595d3e5765ba37b7dff4c652', 'check_answer_text__12b6d7e1507f0e7212c8d9d4764a8a2b', 'check_digit_counts__c2b3326804f7e34493f7827960c36c6d', 'check_targz_has_files__d544b10f933e46c3b8f208241fda6c7e', 'check_file_content__abab8b2b568b2217f8a91073557c53a6', 'check_file_content__2372f00ad964142d3685a78ced2bb21c', 'check_file_format__5d2d411fc8203609cb2ed6eef9424f86', 'check_file_content__65e1f38950e9d333468e23f2dba31d48', 'check_text_alignment__b3837f805dc03cde831accc6a7063290', 'check_file_system__f073ea69e2970c4ed8abb5029189b0ac', 'check_unordered_json_list__d7ef119cb2976be7a75aaacedec7ea6b', 'check_file_contains_function__6c3110d7dd0b28f9da7376f5ce78a46b', 'check_text_alignment__d36eb73e90de91fc69ee1df98b9d5058', 'check_text_color__42394c3fddd5c54d80cb5da5e81a07d8', 'check_sorted_filenames__b150ff4c0a473c93c912a80df923233a', 'check_file_content__05446d05a2730db5bc45314af39c65c5', 'check_line_per_number__2e7f2c4f60f716ea76fd3fbb5849637d', 'check_text_replacement__7061bd4dde7607bb250f2ba7725dc0d5', 'check_text_match__38a0fe7ffb236efbdbbcf3cc7460e807', 'check_output_file__f4c7cbbe6d6bd73cc0b521d770fb3945', 'check_file_count_gte__c0f167f76d6dbaf3d31bd54c91783445', 'check_answer_text__0636655694efccc23299c2ca2f236f1f', 'check_gitlens_and_wordwrap__5c2597160999bc0ba0a4a23d777e3ec1', 'check_file_exists__6e337deb47fb4ded97926f2aaf4149ad_qw35sft2_5e542c11', 'check_file_exists__de94d61f074111bc63d9a066e03aa46f_qw35sft2_d42bca74', 'check_notes_text__dd477a39d59856cb4084c0b8f3f8bce8_qw35sft2_b6600ab9', 'check_file_exists__3ed023e79bfe94e32e5ea286856ad04c_qw35sft2_eee288a4', 'check_all_text_color__2c3419decc4324f20da8cbd8ae0e6c2d_qw35sft2_0c256de4', 'check_stretch_and_text_update__04e2292403d24fd88b4576e3fc170228_qw35sft2_f2180d52', 'check_git_push_and_branch__216e252ae5e5ecc8fbd5efc80755d934_qw35sft2_99e8d888', 'check_path_text_file__a43735a3f0450074dfa18129f92a73a9_qw35sft2_3898dcc4', 'check_file_content__5ca68902218b7e30665245f17696931d_qw35sft2_98d63c75', 'check_invoice_file_saved__d248703d35223088b5508b0c05332ea3_qw35sft2_0cb733ac', 'check_text_file_contains_dimensions__30cadb0ec7eca28df693cf924ab88301_qw35sft2_30caabbd', 'check_shebang_line__7b62591361b18caf2ed6aa916d0c9cf5_qw35sft2_f25870bb', 'check_file_exists__21cbb74de44bf7839b1ea14bdfabea03_qw35sft2_14ee52e3', 'check_file_at_path__82c707062135fc7c567ff519e1b5a575_qw35sft2_1ca24f9b', 'check_killed_and_done_file__aacb4feb23ed663ba8f9729abaff2492_qw35sft2_93f85cfd', 'check_file_exists__057715c74cd9cb2c93d92377c04bb077_qw35sft2_2fb11011', 'check_killed_and_backup_file__2c09d1f072c49dfa84fdfffc2b737545_qw35sft2_968d491a', 'check_file_exists__ea8c7a7a44f5cb78e25fd89d658a863f_qw35sft2_1d22e295', 'check_large_text_no_animations__855f455c8e49beda6b556868bdd9b753_qw35sft2_45a99b91', 'check_all_files_600__1540d35f6307434ff998913b82a97dfa_qw35sft2_067103ba', 'check_user_identity_file__9be4c7481abd67d84bd529b63ed66ca8_qw35sft2_97c039dc', 'check_count_file__5799d2a2e180d6540c452f203d2bde46_qw35sft2_cdeb7860', 'check_large_text_screen_keyboard__0c6c70afb2a250c0e6dcd4a5ca3ca3b8_qw35sft2_a80d73a1', 'check_file_perms_and_rename__653ac9883bb8d908f613f5a6b4d6a405_qw35sft2_f32dae79', 'check_timezone_tokyo__bc8cd5b368955a23d4c5fbd0a520a94a_qw35sft2_3696d128', 'check_sys_info_file__5cefed7a5e164dc00f5c88ae5726c23e_qw35sft2_f6ddcebb', 'check_old_files_cleanup__aba3d770f03c5889eb920c2ccea4f4e0_qw35sft2_55d9accf', 'check_copy_and_create_file2__0f2ab16ddcc8f1a930b6b27f18088a10_qw35sft2_be7a5b0e', 'check_large_text_screen_reader__e5c088a08a33da472617cdfc723aca05_qw35sft2_6141e4a8', 'check_file_and_dir_permissions__6736312e2a9f6c78caa7e5e2f35a4133_qw35sft2_e6e68a10', 'check_python_env_setup__be89f5274d64cd2fc7ba060d990658f8_qw35sft2_78df5e31', 'check_user_audit_files__c74bb543554ec2beed49099c34c9413f_qw35sft2_bd0d0794', 'check_file_org__8f46312161f8fe4e75948b31954b5be0_qw35sft2_c333fad6', 'check_file_copy_3dirs__d110c9f1795d3fe99db0ae18f13e5b66_qw35sft2_06ea6ed3', 'check_timezone_newyork__cc3639ca98d2a9b7b2beb613c1282bfd_qw35sft2_6a4edbdc', 'check_new_files_perms__2e6f4a565bdf8ce8ae9bb4e56e1c79f7_qw35sft2_421bdfed', 'check_home_users_file__baeeb37009528a1d856ed9e685fd5d89_qw35sft2_df5f2a52', 'check_text_scaling_exact__68667bdf1fe16ced5c6ec107d5f22e2d_qw35sft2_92c01ea7', 'check_vim_removed_files_added__fafd25459d3a377f8d0c6d91b93474de_qw35sft2_5c623c89', 'check_rename_and_file__caa3a14501f9d1aa7da658fbd50adb72_qw35sft2_0b8869db', 'check_file_perms_and_archive__00b9eddcbcb085dc89e5e17c3b839024_qw35sft2_ecf3946f', 'check_profile_renamed__c4720772b334cf5b9a4c5019ec6c389a_qw35sft2_d712d113', 'check_new_profile_created__0f2e8c08317dc15425e98e9b4b1171de_qw35sft2_aa3af12e', 'check_default_profile_changed__e7275df5a8446b39f1de20d975873334_qw35sft2_ac1506d0', 'check_file_exists__df188ad33aeda3e315e62cc2c6afb173_qw35sft2_a5bd6840']

def check_file_copy_and_clipboard__267da608b7e8dd72021770568c2da84f(result, expected, **options):
    """Check file was copied and clipboard contains expected path. Partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    expected_path = expected.get('expected_path', '')
    if result.get('file_exists', False):
        score += 0.5
    clipboard = result.get('clipboard', '')
    if expected_path and expected_path in clipboard:
        score += 0.5
    return score

def check_all_text_uppercase__0c637a2f0d91975120dc8e0a76554911(result, expected, **options):
    """Check if all text in the document is uppercase."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    if not paragraphs:
        return 0.0
    total_alpha = 0
    upper_alpha = 0
    for text in paragraphs:
        for ch in text:
            if ch.isalpha():
                total_alpha += 1
                if ch.isupper():
                    upper_alpha += 1
    if total_alpha == 0:
        return 0.0
    ratio = upper_alpha / total_alpha
    if ratio >= 0.95:
        return 1.0
    elif ratio >= 0.8:
        return 0.5
    return 0.0

def check_text_file_keywords__ca2a0f6f5509af90b7b3cf524daf3368(result, expected, **options):
    """Check if text file contains required keywords."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    required_keywords = expected.get('required_keywords', [])
    if not required_keywords:
        return 0.0
    score = 0.0
    total = len(required_keywords)
    for kw in required_keywords:
        if kw.lower() in content.lower():
            score += 1.0 / total
    return min(score, 1.0)

def check_text_replacement__12dc12177d981fb82b80cf3f9509e0cf(result, expected, **options):
    """Check that a find-and-replace was performed correctly.

    Expected rules:
        old_word: str - word that should no longer appear
        new_word: str - word that should appear in its place
        expected_count: int - how many times new_word should appear (at minimum)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else ''
    if not content:
        return 0.0
    old_word = expected.get('old_word', '')
    new_word = expected.get('new_word', '')
    expected_count = expected.get('expected_count', 1)
    if old_word and old_word in content:
        return 0.0
    actual_count = content.count(new_word)
    if actual_count >= expected_count:
        return 1.0
    return 0.0

def check_file_listing__6790a81b7f7c76a950ade3c3fcc6cc45(result, expected, **options):
    """Check if text file contains expected filenames. Partial credit."""
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.3
    expected_names = set(expected.get('expected_names', []))
    if not expected_names:
        return score
    actual_lines = set(result.get('lines', []))
    matches = 0
    for name in expected_names:
        for line in actual_lines:
            if name in line:
                matches += 1
                break
    score += 0.7 * (matches / len(expected_names))
    return min(score, 1.0)

def check_text_content__ec3e3cd160378737f8c3074a7c1fe7f0(result, expected, **options):
    """Check that text file contains expected content."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('content', '').strip()
    expected_value = str(expected.get('expected_value', '')).strip()
    if actual == expected_value:
        return 1.0
    if expected_value and expected_value in actual:
        return 0.5
    return 0.0

def check_outline_view__4b3841fb1241cbb16c8a71b3b365dbfe(result, expected, **options):
    """Check that LibreOffice Impress is in Outline view.

    When Outline view is active (View > Outline), the main editing area changes
    to show a text outline instead of the slide canvas. The accessibility tree
    will contain 'Outline' as part of the view description and the normal slide
    editing elements like 'Click to add Title' will not be present as direct
    editable placeholders.

    Returns 1.0 if in Outline view, 0.0 otherwise.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    tree_str = result if isinstance(result, str) else str(result)
    score = 0.0
    outline_indicators = ['Outline', 'outline']
    for indicator in outline_indicators:
        if indicator in tree_str:
            score += 0.5
            break
    normal_view_indicators = ['Click to add Title', 'Click to add Text']
    normal_found = False
    for indicator in normal_view_indicators:
        if indicator in tree_str:
            normal_found = True
            break
    if not normal_found:
        score += 0.5
    return min(score, 1.0)

def check_download_speed_file__034287361433ba7df6a1a4debdad660e(file_path, expected, **options):
    """Check if download speed file contains valid content with partial credit.

    Scoring:
        0.5 - File exists and has non-empty content
        0.5 - File contains a numeric value representing download speed
    """
    score = 0.0
    if not file_path or not os.path.isfile(str(file_path)):
        logger.debug('check_download_speed_file: file not found or path is empty')
        return 0.0
    try:
        with open(file_path, 'r') as f:
            content = f.read().strip()
    except Exception:
        logger.debug('check_download_speed_file: could not read file')
        return 0.0
    if not content:
        logger.debug('check_download_speed_file: file is empty')
        return 0.0
    score += 0.5
    numbers = re.findall('\\d+\\.?\\d*', content)
    if numbers:
        try:
            val = float(numbers[0])
            if val > 0:
                score += 0.5
        except (ValueError, IndexError):
            pass
    return min(score, 1.0)

def check_python_syntax_valid__94eb7a158102fbe4b2865eae8fd2e0f7(result, expected, **options):
    """Check if the Python file has valid syntax (no SyntaxError)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('valid', False):
        return 1.0
    return 0.0

def check_version_file__4c4a9bd632c3f7487e991bbd8a3f942a(result, expected, **options):
    """Check file exists and contains version pattern. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.4
    content = result.get('content', '')
    app_name = expected.get('app_name', '')
    if app_name and app_name.lower() in content.lower():
        score += 0.3
    version_pattern = expected.get('version_pattern', '\\d+\\.\\d+')
    if re.search(version_pattern, content):
        score += 0.3
    return min(score, 1.0)

def check_line_content__c54a7288908815335d3763c41c9cef77(result, expected, **options):
    """Check that specific lines match expected content. Partial credit per line."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    expected_lines = expected.get('expected_lines', {})
    if not expected_lines:
        return 0.0
    score = 0.0
    total = len(expected_lines)
    for (line_num_str, exp_content) in expected_lines.items():
        line_idx = int(line_num_str) - 1
        if 0 <= line_idx < len(lines) and lines[line_idx] == exp_content:
            score += 1.0 / total
    return min(score, 1.0)

def check_file_renamed__9b65526186f9fb77c93ff9d7de5c5fcb(result, expected, **options):
    """Check if file was renamed from old_name to new_name."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    files = result.get('files', [])
    new_name = expected.get('new_name', '')
    old_name = expected.get('old_name', '')
    score = 0.0
    if new_name in files:
        score += 0.7
    if old_name not in files:
        score += 0.3
    return min(score, 1.0)

def check_file_exists__0ae8da7e0d709c03ee846a1d011df875(result, expected, **options):
    """Check that the expected file exists and has reasonable size."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    min_size = expected.get('min_size', 1000)
    if result.get('file_size', 0) < min_size:
        return 0.5
    return 1.0

def check_line_content__350a72c316de89189f7f590d79389d9e(result, expected, **options):
    """Check that specific lines match expected content. Partial credit per line."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    expected_lines = expected.get('expected_lines', {})
    if not expected_lines:
        return 0.0
    score = 0.0
    total = len(expected_lines)
    for (line_num_str, exp_content) in expected_lines.items():
        line_idx = int(line_num_str) - 1
        if 0 <= line_idx < len(lines) and lines[line_idx] == exp_content:
            score += 1.0 / total
    return min(score, 1.0)

def check_file_organize__a65eb6fd760478532b3e118be8d079f9(result, expected, **options):
    """Check if directory was created and file was moved into it with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir_exists'):
        score += 0.5
    if result.get('file_in_dir'):
        score += 0.5
    return min(score, 1.0)

def check_title_text__1b7b08a28b40e394d9044a5b56ecfe62(result, expected, **options):
    """Check if the slide title matches the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title_text', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if not expected_title:
        return 0.0
    if actual_title.lower() == expected_title.lower():
        return 1.0
    return 0.0

def check_file_moved__e60640f858dbfcce19e818f08302a608(result, expected, **options):
    """Check if file was moved from source to destination. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dst_exists'):
        score += 0.7
    if not result.get('src_exists'):
        score += 0.3
    return min(score, 1.0)

def check_line_content__c8cb9d1d5c3b002c8ddbeb93d3580891(result, expected, **options):
    """Check that specific lines match expected content. Partial credit per line."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    lines = result.get('lines', [])
    expected_lines = expected.get('expected_lines', {})
    if not expected_lines:
        return 0.0
    score = 0.0
    total = len(expected_lines)
    for (line_num_str, exp_content) in expected_lines.items():
        line_idx = int(line_num_str) - 1
        if 0 <= line_idx < len(lines) and lines[line_idx] == exp_content:
            score += 1.0 / total
    return min(score, 1.0)

def check_first_line__35d6e56c2b031e6f068f1907dd34bebf(result, expected, **options):
    """Check if the first line of calculator.py contains the expected comment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    first_line = result.get('first_line', '')
    expected_comment = expected.get('expected_comment', '# Insertion Sort Algorithm')
    if first_line.strip() == expected_comment.strip():
        return 1.0
    if first_line.startswith('#') and 'insertion' in first_line.lower() and ('sort' in first_line.lower()):
        return 0.5
    return 0.0

def check_text_uppercase__0a695088802a9cb1473c98bdb4ec1b4f(result, expected, **options):
    """Check if all text in the document is uppercase."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    text = result.get('text', '')
    if not text:
        return 0.0
    alpha_chars = [c for c in text if c.isalpha()]
    if not alpha_chars:
        return 0.0
    upper_count = sum((1 for c in alpha_chars if c.isupper()))
    ratio = upper_count / len(alpha_chars)
    if ratio >= 0.95:
        return 1.0
    elif ratio >= 0.8:
        return 0.5
    return 0.0

def check_text_alignment__a27e67f69ab078486669e5f21f103d16(result, expected, **options):
    """Check if all paragraph alignments match the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    alignments = result.get('alignments', [])
    expected_alignment = expected.get('expected_alignment', '')
    if not alignments:
        return 0.0
    match_count = sum((1 for a in alignments if a == expected_alignment))
    return match_count / len(alignments)

def check_file_exists_nonempty__29f89216bfdb55e864e9190350cbf19a(result, expected, **options):
    """Check that a file exists and is non-empty.

    Expected keys: min_size (optional), content_contains (optional list of strings).
    Scoring: 0.5 for file existing with content, 0.5 for content keyword match.
    """
    if isinstance(result, str) or not result.get('exists'):
        return 0.0
    score = 0.0
    min_size = expected.get('min_size', 10)
    if result.get('size', 0) >= min_size:
        score += 0.5
    content_contains = expected.get('content_contains', [])
    if not content_contains:
        score += 0.5
    else:
        content_lower = result.get('content', '').lower()
        matched = sum((1 for kw in content_contains if kw.lower() in content_lower))
        if matched > 0:
            score += 0.5 * (matched / len(content_contains))
    return min(score, 1.0)

def check_json_settings_subset__7efbe3197e1cdd53b30f5426c6d455b2(actual, expected, **options):
    """Check if all expected settings exist in actual VS Code settings using deep subset matching.

    For nested dict values (e.g. files.exclude), checks that all expected
    key-value pairs are present in the actual dict, rather than requiring
    exact equality. This handles VS Code writing default patterns alongside
    user-added entries.

    Args:
        actual (str): path to the settings.json file
        expected (dict): expected dict containing key "expected" with settings to verify

    Returns:
        float: 1.0 if all expected settings found, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    if not expect:
        return 0.0
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel:
            return 0.0
        if isinstance(value, dict) and isinstance(actual_value, dict):
            for (sub_key, sub_val) in value.items():
                if actual_value.get(sub_key) != sub_val:
                    return 0.0
        elif actual_value != value:
            return 0.0
    return 1.0

def check_text_scaling__04c87fb26c99f62f08868d28d96db622(result, expected, **options):
    """Check if text scaling factor matches the expected value.

    Args:
        result: Output from vm_command_line (gsettings get text-scaling-factor).
        expected: Dict with 'expected_factor' key (rules dict passed directly).
    """
    try:
        if isinstance(result, dict) and result.get('error'):
            return 0.0
        scaling_str = result if isinstance(result, str) else str(result)
        scaling_str = scaling_str.strip()
        scaling_factor = float(scaling_str)
        expected_factor = float(expected.get('expected_factor', 1.25))
        return 1.0 if abs(scaling_factor - expected_factor) < 0.01 else 0.0
    except (ValueError, TypeError, AttributeError):
        return 0.0

def check_direct_json_object(result, rules) -> float:
    """
    One of the most commonly used function to evalute.
    Compare two json objects directly.

    Fixed version that correctly handles 'ignore_list_order' parameter
    from the top-level rules dict instead of expected_json dict.
    """
    logger.info(f'[DEBUG] check_direct_json_object called with result: {result}')
    logger.info(f'[DEBUG] check_direct_json_object called with rules: {rules}')
    if isinstance(result, str):
        result = result.strip()
        result = result.replace("'", '"')
        result = json.loads(result)
    logger.info(f'[DEBUG] Processed result: {result}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    try:
        expected_json = rules.get('expected', {})
        if expected_json:
            for (key, value) in expected_json.items():
                if value == '__EVALUATION_FAILED__':
                    logger.error(f"[DEBUG] Expected value for key '{key}' indicates evaluation failure, returning 0.0")
                    return 0.0
    except Exception as e:
        logger.error(f'[DEBUG] Error checking for evaluation failure indicator: {e}')
        return 0.0
    try:
        expect_in_result = rules.get('expect_in_result', False)
        logger.info(f'[DEBUG] expect_in_result: {expect_in_result}')
        if not expect_in_result:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON: {expected_json}')
            for key in expected_json.keys():
                expected_value = expected_json.get(key)
                actual_value = result.get(key)
                logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
                if rules.get('ignore_list_order', False):
                    expected_value = sorted(expected_value)
                    result_value = sorted(result.get(key))
                    logger.info(f'[DEBUG] Comparing lists (sorted): expected={expected_value}, actual={result_value}')
                    if expected_value != result_value:
                        logger.info(f"[DEBUG] List comparison failed for key '{key}', returning 0.0")
                        return 0.0
                elif expected_value != actual_value:
                    logger.info(f"[DEBUG] Value comparison failed for key '{key}': expected='{expected_value}', actual='{actual_value}', returning 0.0")
                    return 0.0
                else:
                    logger.info(f"[DEBUG] Value comparison passed for key '{key}'")
            logger.info('[DEBUG] All comparisons passed, returning 1.0')
            return 1.0
        else:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON (expect_in_result mode): {expected_json}')
            for key in expected_json.keys():
                if isinstance(expected_json.get(key), list):
                    flag = 0
                    expected_value_list = expected_json.get(key)
                    logger.info(f"[DEBUG] Checking list key '{key}': expected_list={expected_value_list}, actual='{result.get(key)}'")
                    for each_expected_value in expected_value_list:
                        if isinstance(result.get(key), list) and each_expected_value in result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' in result list for key '{key}'")
                            break
                        elif isinstance(result.get(key), str) and each_expected_value == result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' matches result string for key '{key}'")
                            break
                    if flag == 0:
                        logger.info(f"[DEBUG] No expected values found in result for key '{key}', returning 0.0")
                        return 0.0
                elif isinstance(expected_json.get(key), str):
                    expected_str = expected_json.get(key)
                    actual_str = result.get(key)
                    logger.info(f"[DEBUG] Checking string key '{key}': expected='{expected_str}', actual='{actual_str}'")
                    if expected_str not in actual_str:
                        logger.info(f"[DEBUG] Expected string '{expected_str}' not found in actual string '{actual_str}' for key '{key}', returning 0.0")
                        return 0.0
                else:
                    logger.debug('check_direct_json_object: expected value type not supported')
                    return 0.0
            logger.info('[DEBUG] All expect_in_result comparisons passed, returning 1.0')
            return 1.0
    except Exception as e:
        logger.debug(f'check_direct_json_object: result is not a valid json object, error: {e}')
        return 0.0

def check_file_dual_location__454f85dd5a1086d0f2418681bdac8477(result, expected, **options):
    """Check if file exists at both locations with partial credit.

    Scoring:
    - 0.5: File exists at location A (Desktop - restored from Trash)
    - 0.5: File exists at location B (Documents - copy)
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('location_a', False):
        score += 0.5
    if result.get('location_b', False):
        score += 0.5
    return min(score, 1.0)

def check_file_content__c80ddc5f2ae016c0c4e6f2749cc5143e(result, expected, **options):
    """Check if file content matches expected content with partial credit."""
    if not result:
        return 0.0
    actual = result.strip() if isinstance(result, str) else str(result).strip()
    expected_content = expected.get('expected_content', '').strip()
    if actual == expected_content:
        return 1.0
    actual_lines = [l.strip() for l in actual.split('\n') if l.strip()]
    expected_lines = [l.strip() for l in expected_content.split('\n') if l.strip()]
    if not expected_lines:
        return 0.0
    correct = sum((1 for (a, e) in zip(actual_lines, expected_lines) if a == e))
    return correct / max(len(expected_lines), len(actual_lines))

def check_file_rename__5128f87c595d3e5765ba37b7dff4c652(result, expected, **options):
    """Check if file rename was successful with partial credit.

    Scoring:
    - 0.5: New file name exists at expected location
    - 0.5: Old file name no longer exists (was renamed, not copied)
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('new_exists', False):
        score += 0.5
    if result.get('old_gone', False):
        score += 0.5
    return min(score, 1.0)

def check_answer_text__12b6d7e1507f0e7212c8d9d4764a8a2b(result, expected, **options):
    """Check if the answer text matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('answer_text', '')
    if actual is None:
        return 0.0
    expected_answer = expected.get('expected_answer', '')
    actual_clean = actual.strip().lower()
    expected_clean = expected_answer.strip().lower()
    if actual_clean == expected_clean:
        return 1.0
    if len(expected_clean) > 0 and len(actual_clean) > 0:
        matches = sum((1 for (a, e) in zip(actual_clean, expected_clean) if a == e))
        max_len = max(len(actual_clean), len(expected_clean))
        return matches / max_len
    return 0.0

def check_digit_counts__c2b3326804f7e34493f7827960c36c6d(result, expected, **options):
    """Check column E header and digit count values."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header')
    expected_header = expected.get('expected_header', 'Digit Count')
    if header and expected_header.lower() in str(header).lower():
        score += 0.1
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return score
    total = len(expected_values)
    matches = 0
    for i in range(min(len(actual_values), total)):
        actual = actual_values[i]
        exp = expected_values[i]
        if actual is None:
            continue
        try:
            if int(float(actual)) == exp:
                matches += 1
        except (ValueError, TypeError):
            continue
    if total > 0:
        score += 0.9 * (matches / total)
    return min(score, 1.0)

def check_targz_has_files__d544b10f933e46c3b8f208241fda6c7e(result, expected, **options):
    """Check if the tar.gz archive contains the expected files with partial credit."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    required_extensions = expected.get('required_extensions', [])
    files = result.get('files', [])
    files_lower = [f.lower() for f in files]
    if not required_extensions:
        return 1.0 if len(files) > 0 else 0.0
    score = 0.0
    per_file = 1.0 / len(required_extensions) if required_extensions else 0.0
    for ext in required_extensions:
        ext_lower = ext.lower()
        for f in files_lower:
            if f.endswith(ext_lower):
                score += per_file
                break
    return min(score, 1.0)

def check_file_content__abab8b2b568b2217f8a91073557c53a6(result, expected, **options):
    """Check if file content matches expected content with partial credit."""
    if not result:
        return 0.0
    actual = result.strip() if isinstance(result, str) else str(result).strip()
    expected_content = expected.get('expected_content', '').strip()
    if actual == expected_content:
        return 1.0
    actual_lines = [l.strip() for l in actual.split('\n') if l.strip()]
    expected_lines = [l.strip() for l in expected_content.split('\n') if l.strip()]
    if not expected_lines:
        return 0.0
    correct = sum((1 for (a, e) in zip(actual_lines, expected_lines) if a == e))
    return correct / max(len(expected_lines), len(actual_lines))

def check_file_content__2372f00ad964142d3685a78ced2bb21c(result, expected, **options):
    """Check file content contains required strings with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    required_strings = expected.get('required_strings', [])
    if not required_strings:
        return 0.0
    score_per_item = 1.0 / len(required_strings)
    total_score = 0.0
    for s in required_strings:
        if s.lower() in content.lower():
            total_score += score_per_item
    return min(total_score, 1.0)

def check_file_format__5d2d411fc8203609cb2ed6eef9424f86(result, expected, **options):
    """Check if file exists and matches expected format."""
    if isinstance(result, dict) and (not result.get('exists', False)):
        return 0.0
    score = 0.0
    score += 0.5
    expected_format = expected.get('format', 'JPEG')
    actual_format = result.get('format', '')
    if actual_format and actual_format.upper() == expected_format.upper():
        score += 0.5
    return score

def check_file_content__65e1f38950e9d333468e23f2dba31d48(result, expected, **options):
    """Check if file content matches expected content with partial credit."""
    if not result:
        return 0.0
    actual = result.strip() if isinstance(result, str) else str(result).strip()
    expected_content = expected.get('expected_content', '').strip()
    if actual == expected_content:
        return 1.0
    actual_lines = [l.strip() for l in actual.split('\n') if l.strip()]
    expected_lines = [l.strip() for l in expected_content.split('\n') if l.strip()]
    if not expected_lines:
        return 0.0
    correct = sum((1 for (a, e) in zip(actual_lines, expected_lines) if a == e))
    return correct / max(len(expected_lines), len(actual_lines))

def check_text_alignment__b3837f805dc03cde831accc6a7063290(result, expected, **options):
    """Check if all paragraph alignments match the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    alignments = result.get('alignments', [])
    expected_alignment = expected.get('expected_alignment', '')
    if not alignments:
        return 0.0
    match_count = sum((1 for a in alignments if a == expected_alignment))
    return match_count / len(alignments)

def check_file_system__f073ea69e2970c4ed8abb5029189b0ac(result, expected, **options):
    """Partial credit: 0.5 for directory existing, 0.5 for file existing."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir_exists'):
        score += 0.5
    if result.get('file_exists'):
        score += 0.5
    return score

def check_unordered_json_list__d7ef119cb2976be7a75aaacedec7ea6b(result, expected, **options):
    """
    Compare JSON objects with order-independent list comparison.
    For each key in expected, if the value is a list, compare sorted versions.
    Otherwise compare directly.
    """
    if result is None:
        return 0.0
    try:
        expected_json = expected.get('expected', {})
        if not expected_json:
            return 0.0
        for (key, expected_value) in expected_json.items():
            actual_value = result.get(key)
            if actual_value is None:
                logger.info(f"Key '{key}' not found in result, returning 0.0")
                return 0.0
            if isinstance(expected_value, list) and isinstance(actual_value, list):
                if sorted(expected_value) != sorted(actual_value):
                    logger.info(f"List mismatch for key '{key}': expected={sorted(expected_value)}, actual={sorted(actual_value)}")
                    return 0.0
            elif expected_value != actual_value:
                logger.info(f"Value mismatch for key '{key}': expected={expected_value}, actual={actual_value}")
                return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error in check_unordered_json_list: {e}')
        return 0.0

def check_file_contains_function__6c3110d7dd0b28f9da7376f5ce78a46b(result, expected, **options):
    """Check if file contains expected function definition with partial credit."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '')
    score = 0.0
    if content.strip():
        score += 0.5
    func_name = expected.get('function_name', '')
    if func_name and 'def ' + func_name in content:
        score += 0.5
    return min(score, 1.0)

def check_text_alignment__d36eb73e90de91fc69ee1df98b9d5058(result, expected, **options):
    """Check if all paragraph alignments match the expected alignment."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    alignments = result.get('alignments', [])
    expected_alignment = expected.get('expected_alignment', '')
    if not alignments:
        return 0.0
    match_count = sum((1 for a in alignments if a == expected_alignment))
    return match_count / len(alignments)

def check_text_color__42394c3fddd5c54d80cb5da5e81a07d8(result, expected, **options):
    """Check if text color matches expected color with tolerance."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_color = result.get('color', '').upper()
    expected_color = expected.get('expected_color', '').upper()
    if not actual_color or not expected_color:
        return 0.0
    if actual_color == expected_color:
        return 1.0
    try:
        ar = int(actual_color[0:2], 16)
        ag = int(actual_color[2:4], 16)
        ab = int(actual_color[4:6], 16)
        er = int(expected_color[0:2], 16)
        eg = int(expected_color[2:4], 16)
        eb = int(expected_color[4:6], 16)
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        threshold = expected.get('tolerance', 30)
        if distance <= threshold:
            return 1.0
    except (ValueError, IndexError):
        pass
    return 0.0

def check_sorted_filenames__b150ff4c0a473c93c912a80df923233a(result, expected, **options):
    """Check if the text file contains the expected sorted filenames.
    Partial credit: each correct filename in correct position = 0.25.
    """
    if result.get('error'):
        return 0.0
    actual_lines = result.get('lines', [])
    expected_lines = expected.get('expected_lines', [])
    if not expected_lines:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_lines)
    for (i, exp) in enumerate(expected_lines):
        if i < len(actual_lines):
            actual_clean = actual_lines[i].strip().lower()
            exp_clean = exp.strip().lower()
            if exp_clean == actual_clean or exp_clean in actual_clean:
                score += per_item
    return min(score, 1.0)

def check_file_content__05446d05a2730db5bc45314af39c65c5(result, expected, **options):
    """Check if file content matches expected value."""
    if result.get('error'):
        return 0.0
    actual = result.get('content', '').strip()
    expected_val = str(expected.get('expected_content', '')).strip()
    if actual == expected_val:
        return 1.0
    try:
        if float(actual) == float(expected_val):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_line_per_number__2e7f2c4f60f716ea76fd3fbb5849637d(result, expected, **options):
    """Check if res.txt contains each sorted number on its own line without prefix.

    Args:
        result: local file path to res.txt (from vm_file getter)
        expected: dict with 'expected_lines' list of string numbers
    Returns:
        float: 1.0 if format and values match, partial credit for correct values wrong format
    """
    if not result:
        return 0.0
    try:
        with open(result) as f:
            content = f.read().strip()
    except Exception:
        return 0.0
    expected_lines = expected.get('expected_lines', [])
    actual_lines = [line.strip() for line in content.split('\n') if line.strip()]
    if actual_lines == expected_lines:
        return 1.0
    import re
    actual_numbers = [int(x) for x in re.findall('\\d+', content)]
    expected_numbers = [int(x) for x in expected_lines]
    if actual_numbers == expected_numbers:
        return 0.5
    return 0.0

def check_text_replacement__7061bd4dde7607bb250f2ba7725dc0d5(result, expected, **options):
    """Check that a find-and-replace was performed correctly.

    Expected rules:
        old_word: str - word that should no longer appear
        new_word: str - word that should appear in its place
        expected_count: int - how many times new_word should appear (at minimum)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    content = result.get('content', '') if isinstance(result, dict) else ''
    if not content:
        return 0.0
    old_word = expected.get('old_word', '')
    new_word = expected.get('new_word', '')
    expected_count = expected.get('expected_count', 1)
    if old_word and old_word in content:
        return 0.0
    actual_count = content.count(new_word)
    if actual_count >= expected_count:
        return 1.0
    return 0.0

def check_text_match__38a0fe7ffb236efbdbbcf3cc7460e807(result, expected, **options):
    """Check if text from a PPTX shape matches the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_text = result.get('text', '').strip().upper()
    expected_text = expected.get('expected_text', '').strip().upper()
    if actual_text == expected_text:
        return 1.0
    if expected_text and (expected_text in actual_text or actual_text in expected_text):
        return 0.5
    return 0.0

def check_output_file__f4c7cbbe6d6bd73cc0b521d770fb3945(result, expected, **options):
    """Check if output.txt exists and contains expected number of output lines."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.5
    lines = result.get('lines', [])
    expected_count = expected.get('expected_line_count', 5)
    if len(lines) == expected_count:
        all_numeric = all((l.lstrip('-').isdigit() for l in lines))
        if all_numeric:
            score += 0.5
    return min(score, 1.0)

def check_file_count_gte__c0f167f76d6dbaf3d31bd54c91783445(result, expected, **options):
    """Check if the file count from vm_command_line output is >= min_count.

    Args:
        result: String output from vm_command_line (e.g., "1
")
        expected: Dict with 'min_count' key (default 1)
    Returns:
        1.0 if count >= min_count, else 0.0
    """
    if not result:
        return 0.0
    try:
        count = int(result.strip())
        min_count = expected.get('min_count', 1)
        return 1.0 if count >= min_count else 0.0
    except (ValueError, TypeError):
        return 0.0

def check_answer_text__0636655694efccc23299c2ca2f236f1f(result, expected, **options):
    """Check if the answer text matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('answer_text', '')
    if actual is None:
        return 0.0
    expected_answer = expected.get('expected_answer', '')
    actual_clean = actual.strip().lower()
    expected_clean = expected_answer.strip().lower()
    if actual_clean == expected_clean:
        return 1.0
    if len(expected_clean) > 0 and len(actual_clean) > 0:
        matches = sum((1 for (a, e) in zip(actual_clean, expected_clean) if a == e))
        max_len = max(len(actual_clean), len(expected_clean))
        return matches / max_len
    return 0.0

def check_gitlens_and_wordwrap__5c2597160999bc0ba0a4a23d777e3ec1(result, expected, **options):
    """Check if GitLens is installed and word wrap is enabled. Partial credit: 0.5 each."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    extension_id = expected.get('extension_id', 'eamodio.gitlens')
    extensions = result.get('extensions', '')
    if extension_id.lower() in extensions.lower():
        score += 0.5
    setting_key = expected.get('setting_key', 'editor.wordWrap')
    setting_value = expected.get('setting_value', 'on')
    settings = result.get('settings', {})
    actual_value = settings.get(setting_key)
    if str(actual_value) == str(setting_value):
        score += 0.5
    return min(score, 1.0)

def check_file_exists__6e337deb47fb4ded97926f2aaf4149ad_qw35sft2_5e542c11(result, expected, **options):
    """Check that the exported file exists and is a valid Windows BMP (magic bytes 'BM')."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    if not result.get('exists'):
        return 0.0
    header_bytes = result.get('header_bytes', b'')
    if header_bytes[:2] != b'BM':
        return 0.0
    return 1.0

def check_file_exists__de94d61f074111bc63d9a066e03aa46f_qw35sft2_d42bca74(result, expected, **options):
    """
    Partial-credit check for animated GIF conversion from the first 2 seconds of video.

    Scoring:
      0.4 — file exists at /home/user/clip.gif AND is a valid GIF (GIF87a/GIF89a)
      0.3 — GIF is animated (frame_count > 1)
      0.3 — GIF duration is ~2 seconds (1.0–3.5 s); if delays are unavailable,
             a plausible frame count (5–120) is accepted as a proxy
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists') and result.get('is_gif'):
        score += 0.4
    else:
        return 0.0
    frame_count = result.get('frame_count', 0)
    if frame_count > 1:
        score += 0.3
    duration = result.get('duration_secs', 0.0)
    if 1.0 <= duration <= 3.5:
        score += 0.3
    elif duration == 0.0 and 5 <= frame_count <= 120:
        score += 0.3
    return min(score, 1.0)

def check_notes_text__dd477a39d59856cb4084c0b8f3f8bce8_qw35sft2_b6600ab9(result, expected, **options):
    """Check that slide 2 notes text matches expected. Case-insensitive strip comparison."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    notes_text = result.get('notes_text', '')
    expected_notes = expected.get('expected_notes', '')
    if notes_text.strip().lower() == expected_notes.strip().lower():
        return 1.0
    return 0.0

def check_file_exists__3ed023e79bfe94e32e5ea286856ad04c_qw35sft2_eee288a4(result, expected, **options):
    """Return 1.0 if the expected file was created, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('exists', False) else 0.0

def check_all_text_color__2c3419decc4324f20da8cbd8ae0e6c2d_qw35sft2_0c256de4(result, expected, **options):
    """Check dark red color on title (0.34), body (0.33), and table (0.33) with partial credit."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_color_ok') is True:
        score += 0.34
    if result.get('body_color_ok') is True:
        score += 0.33
    if result.get('table_color_ok') is True:
        score += 0.33
    return round(min(score, 1.0), 2)

def check_stretch_and_text_update__04e2292403d24fd88b4576e3fc170228_qw35sft2_f2180d52(result, expected, **options):
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tolerance_cm = 0.8
    slide_width = result.get('slide_width_cm', 25.4)
    slide_height = result.get('slide_height_cm', 19.05)
    image_width = result.get('image_width_cm', 0.0)
    image_height = result.get('image_height_cm', 0.0)
    image_left = result.get('image_left_cm', -99.0)
    image_top = result.get('image_top_cm', -99.0)
    fills_width = abs(image_width - slide_width) <= tolerance_cm
    fills_height = abs(image_height - slide_height) <= tolerance_cm
    image_fills_page = fills_width or fills_height
    if fills_width:
        expected_top = (slide_height - image_height) / 2
        centered = abs(image_left) <= tolerance_cm and abs(image_top - expected_top) <= tolerance_cm
    elif fills_height:
        expected_left = (slide_width - image_width) / 2
        centered = abs(image_top) <= tolerance_cm and abs(image_left - expected_left) <= tolerance_cm
    else:
        centered = False
    if image_fills_page and centered:
        score += 0.5
    expected_text = expected.get('expected_text', 'Background Investigation Process')
    actual_text = result.get('textbox9_text') or result.get('bottom_text_fallback') or ''
    actual_normalized = ' '.join(actual_text.split()).strip()
    expected_normalized = ' '.join(expected_text.split()).strip()
    if actual_normalized == expected_normalized:
        score += 0.5
    return min(score, 1.0)

def check_git_push_and_branch__216e252ae5e5ecc8fbd5efc80755d934_qw35sft2_99e8d888(result, expected, **options):
    """
    Partial-credit metric checking:
      0.5 — Remote repo's latest commit message contains expected commit message
      0.5 — Local binder repo has the expected branch name
    Returns float in [0.0, 1.0].
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_msg = expected.get('commit_message', 'daily update')
    remote_log = result.get('remote_log', '')
    if expected_msg.lower() in remote_log.lower():
        score += 0.5
    expected_branch = expected.get('branch_name', 'backup')
    branches = result.get('branches', '')
    branch_list = [b.strip().lstrip('* ').strip() for b in branches.splitlines() if b.strip()]
    if expected_branch in branch_list:
        score += 0.5
    return min(score, 1.0)

def check_path_text_file__a43735a3f0450074dfa18129f92a73a9_qw35sft2_3898dcc4(result, expected, **options):
    """Check that the text file contains exactly the expected file path."""
    if result is None or result.get('error'):
        return 0.0
    content = result.get('content', '').strip()
    expected_path = expected.get('expected_path', '').strip()
    if not expected_path:
        return 0.0
    if content == expected_path:
        return 1.0
    first_line = content.splitlines()[0].strip() if content else ''
    if first_line == expected_path:
        return 0.8
    return 0.0

def check_file_content__5ca68902218b7e30665245f17696931d_qw35sft2_98d63c75(result, expected, **options):
    """Check that ~/Test/Speed/results.txt contains an expected line."""
    if not isinstance(result, dict) or result.get('error') or (not result.get('lines')):
        return 0.0
    expected_line = expected.get('expected_line', '')
    if not expected_line:
        return 1.0 if result.get('lines') else 0.0
    actual_lines = result.get('lines', [])
    for line in actual_lines:
        if line.strip() == expected_line.strip():
            return 1.0
    return 0.0

def check_invoice_file_saved__d248703d35223088b5508b0c05332ea3_qw35sft2_0cb733ac(result, expected, **options):
    """Check that the December AWS invoice was saved with the correct naming pattern."""
    if not result:
        return 0.0
    expected_filename = expected.get('expected_filename', 'aws-invoice-2312.pdf')
    files = [f.strip() for f in result.strip().splitlines() if f.strip()]
    return 1.0 if expected_filename in files else 0.0

def check_text_file_contains_dimensions__30cadb0ec7eca28df693cf924ab88301_qw35sft2_30caabbd(result, expected, **options):
    """Check if a text file contains both the expected width and height values."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    content = result.get('content', '') or ''
    width = str(expected.get('expected_width', '320'))
    height = str(expected.get('expected_height', '510'))
    score = 0.0
    if width in content:
        score += 0.5
    if height in content:
        score += 0.5
    return score

def check_shebang_line__7b62591361b18caf2ed6aa916d0c9cf5_qw35sft2_f25870bb(result, expected, **options):
    """Check that the first line of main.py is the expected shebang.

    Returns 1.0 if the first line matches '#!/usr/bin/env python3', 0.0 otherwise.
    """
    if not isinstance(result, str):
        return 0.0
    expected_shebang = expected.get('shebang', '#!/usr/bin/env python3')
    return 1.0 if result.strip() == expected_shebang else 0.0

def check_file_exists__21cbb74de44bf7839b1ea14bdfabea03_qw35sft2_14ee52e3(result, expected, **options):
    """Return 1.0 if the file exists, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('exists') else 0.0

def check_file_at_path__82c707062135fc7c567ff519e1b5a575_qw35sft2_1ca24f9b(result, expected, **options):
    """Check that the target file exists on the filesystem."""
    if result is None or result.get('error'):
        return 0.0
    return 1.0 if result.get('exists', False) else 0.0

def check_killed_and_done_file__aacb4feb23ed663ba8f9729abaff2492_qw35sft2_93f85cfd(result, expected, **options):
    """
    Check that LibreOffice was killed (0.5) AND /home/user/Desktop/done.txt
    contains the expected text (0.5).

    result: stdout from composite shell command with 'proc_status:...|file_content:...'
    expected: already-unwrapped rules dict with 'expected_text' key
    """
    if result is None:
        return 0.0
    result_str = str(result)
    score = 0.0
    if 'proc_status:not_running' in result_str:
        score += 0.5
    expected_text = expected.get('expected_text', 'terminated')
    try:
        after_marker = result_str.split('file_content:', 1)[1]
        file_content = after_marker.strip()
        if expected_text.lower() in file_content.lower():
            score += 0.5
    except (IndexError, AttributeError):
        pass
    return min(score, 1.0)

def check_file_exists__057715c74cd9cb2c93d92377c04bb077_qw35sft2_2fb11011(result, expected, **options):
    """Check if the file existence matches expected state."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_exists = expected.get('expected_exists', True)
    return 1.0 if result.get('exists') == expected_exists else 0.0

def check_killed_and_backup_file__2c09d1f072c49dfa84fdfffc2b737545_qw35sft2_968d491a(result, expected, **options):
    """
    Check that LibreOffice was killed (0.5) AND the document has been renamed
    to backup.docx on the Desktop (0.5).

    result: stdout from composite shell command with 'proc_status:...|backup_exists:...'
    expected: already-unwrapped rules dict
    """
    if result is None:
        return 0.0
    result_str = str(result)
    score = 0.0
    if 'proc_status:not_running' in result_str:
        score += 0.5
    if 'backup_exists:yes' in result_str:
        score += 0.5
    return min(score, 1.0)

def check_file_exists__ea8c7a7a44f5cb78e25fd89d658a863f_qw35sft2_1d22e295(result, expected, **options):
    """Check if a file exists with minimum required size and correct duration."""
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    min_size = expected.get('min_size', 1)
    if result.get('size', 0) < min_size:
        return 0.0
    min_duration = expected.get('min_duration', None)
    max_duration = expected.get('max_duration', None)
    if min_duration is not None and max_duration is not None:
        duration = result.get('duration', -1.0)
        if duration < 0 or not min_duration <= duration <= max_duration:
            return 0.0
    return 1.0

def check_large_text_no_animations__855f455c8e49beda6b556868bdd9b753_qw35sft2_45a99b91(result, expected, **options):
    """Check large-text enabled and animations disabled, with partial credit (0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('large_text') is True:
        score += 0.5
    if result.get('animations_enabled') is False:
        score += 0.5
    return score

def check_all_files_600__1540d35f6307434ff998913b82a97dfa_qw35sft2_067103ba(result, expected, **options):
    """Check that all regular files in testDir have permission 600."""
    if not result or 'error' in str(result).lower():
        return 0.0
    permissions = result.get('permissions', {})
    if not permissions:
        return 0.0
    expected_perm = expected.get('expected_permission', '600')
    all_correct = all((perm == expected_perm for perm in permissions.values()))
    return 1.0 if all_correct else 0.0

def check_user_identity_file__9be4c7481abd67d84bd529b63ed66ca8_qw35sft2_97c039dc(result, expected, **options):
    """Check that ~/user_identity.txt exists and contains uid/user info from `id`."""
    if result.get('error') or not result.get('content'):
        return 0.0
    content = result['content'].lower()
    keywords = expected.get('keywords', [])
    if all((kw.lower() in content for kw in keywords)):
        return 1.0
    return 0.0

def check_count_file__5799d2a2e180d6540c452f203d2bde46_qw35sft2_cdeb7860(result, expected, **options):
    """Partial credit: 0.34 archive with correct old files, 0.33 new_files populated, 0.33 count.txt correct."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('archive_exists') and result.get('archive_has_old_files'):
        score += 0.34
    if result.get('new_files_populated'):
        score += 0.33
    expected_count = str(expected.get('expected_count', '2'))
    if result.get('count_content', '').strip() == expected_count:
        score += 0.33
    return round(score, 2)

def check_large_text_screen_keyboard__0c6c70afb2a250c0e6dcd4a5ca3ca3b8_qw35sft2_a80d73a1(result, expected, **options):
    """Check large-text enabled and screen keyboard enabled, with partial credit (0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('large_text') is True:
        score += 0.5
    if result.get('screen_keyboard') is True:
        score += 0.5
    return score

def check_file_perms_and_rename__653ac9883bb8d908f613f5a6b4d6a405_qw35sft2_f32dae79(result, expected, **options):
    """Check all files are 644 (0.5 pts) and file3.txt was renamed to file3_backup.txt (0.5 pts)."""
    if not result:
        return 0.0
    score = 0.0
    file_perms = result.get('file_permissions', [])
    expected_perm = expected.get('expected_file_permission', '644')
    if file_perms and all((p == expected_perm for p in file_perms)):
        score += 0.5
    if result.get('backup_file_exists', False) and result.get('old_file_gone', False):
        score += 0.5
    return score

def check_timezone_tokyo__bc8cd5b368955a23d4c5fbd0a520a94a_qw35sft2_3696d128(result, expected, **options):
    """Check if system timezone is set to Asia/Tokyo (Japan Standard Time).

    Parses `timedatectl status` output (line index 3 is the Time zone line).
    Example line: '                Time zone: Asia/Tokyo (JST, +0900)'
    """
    if result is None:
        return 0.0
    try:
        lines = str(result).strip().split('\n')
        if len(lines) < 4:
            return 0.0
        tz_line = lines[3]
        colon_idx = tz_line.find(':')
        if colon_idx == -1:
            return 0.0
        tz_part = tz_line[colon_idx + 1:].strip()
        tz_name = tz_part.split(' ')[0]
        return 1.0 if 'Tokyo' in tz_name else 0.0
    except Exception:
        return 0.0

def check_sys_info_file__5cefed7a5e164dc00f5c88ae5726c23e_qw35sft2_f6ddcebb(result, expected, **options):
    """Check ~/sys_info.txt contains both username and hostname strings."""
    if result.get('error') or not result.get('content'):
        return 0.0
    content = result['content']
    include_terms = expected.get('include', [])
    if all((term in content for term in include_terms)):
        return 1.0
    return 0.0

def check_old_files_cleanup__aba3d770f03c5889eb920c2ccea4f4e0_qw35sft2_55d9accf(result, expected, **options):
    """Partial credit: 0.33 archive exists in old_files, 0.34 no loose txt files in old_files root, 0.33 recent txt files present in new_files."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('archive_exists'):
        score += 0.33
    if result.get('loose_txt_count', 1) == 0:
        score += 0.34
    if result.get('new_txt_count', 0) > 0:
        score += 0.33
    return round(min(score, 1.0), 2)

def check_copy_and_create_file2__0f2ab16ddcc8f1a930b6b27f18088a10_qw35sft2_be7a5b0e(result, expected, **options):
    """0.25 each: file1 in dir1, dir2, dir3; file2 exists in home dir."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir1_has_file1'):
        score += 0.25
    if result.get('dir2_has_file1'):
        score += 0.25
    if result.get('dir3_has_file1'):
        score += 0.25
    if result.get('home_has_file2'):
        score += 0.25
    return min(round(score, 6), 1.0)

def check_large_text_screen_reader__e5c088a08a33da472617cdfc723aca05_qw35sft2_6141e4a8(result, expected, **options):
    """Check large-text enabled and screen-reader enabled, with partial credit (0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('large_text') is True:
        score += 0.5
    if result.get('screen_reader') is True:
        score += 0.5
    return score

def check_file_and_dir_permissions__6736312e2a9f6c78caa7e5e2f35a4133_qw35sft2_e6e68a10(result, expected, **options):
    """Check files are 644 (0.5 pts) and subdirectories are 755 (0.5 pts)."""
    if not result:
        return 0.0
    score = 0.0
    file_perms = result.get('file_permissions', [])
    dir_perms = result.get('dir_permissions', [])
    expected_file = expected.get('expected_file_permission', '644')
    expected_dir = expected.get('expected_dir_permission', '755')
    if file_perms and all((p == expected_file for p in file_perms)):
        score += 0.5
    if dir_perms and all((p == expected_dir for p in dir_perms)):
        score += 0.5
    return score

def check_python_env_setup__be89f5274d64cd2fc7ba060d990658f8_qw35sft2_78df5e31(result, expected, **options):
    """
    Partial-credit check for two Python environment setup goals:
      - 0.5 if PYTHONUTF8=1 is present in /etc/environment
      - 0.5 if ~/python_info.txt exists and contains Python 3.x version info
    """
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    if 'PYTHONUTF8=1' in result.get('env_content', ''):
        score += 0.5
    file_content = result.get('file_content', '')
    if 'Python 3' in file_content or '3.10' in file_content:
        score += 0.5
    return score

def check_user_audit_files__c74bb543554ec2beed49099c34c9413f_qw35sft2_bd0d0794(result, expected, **options):
    """Partial-credit check: 0.5 for correct whoami.txt, 0.5 for correct users.txt."""
    score = 0.0
    whoami_expected = expected.get('whoami_expected', 'user')
    users_keyword = expected.get('users_keyword', 'user:')
    whoami_content = result.get('whoami_content', '').strip()
    users_content = result.get('users_content', '')
    if whoami_expected.lower() in whoami_content.lower():
        score += 0.5
    if users_keyword in users_content:
        score += 0.5
    return score

def check_file_org__8f46312161f8fe4e75948b31954b5be0_qw35sft2_c333fad6(result, expected, **options):
    """
    Partial credit scoring:
    0.5 - A .tar.gz archive exists in /tmp/test_files/old_files/ AND
          it contains at least 2 .txt files (old files were actually compressed into it).
    0.5 - /tmp/test_files/new_files/ contains at least 2 .txt files
          (recently modified files were moved there by the agent).
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('archive_found') and result.get('archive_has_txt_files'):
        score += 0.5
    min_count = expected.get('new_files_min_count', 2)
    if result.get('new_files_count', 0) >= min_count:
        score += 0.5
    return round(score, 2)

def check_file_copy_3dirs__d110c9f1795d3fe99db0ae18f13e5b66_qw35sft2_06ea6ed3(result, expected, **options):
    """Award 1/3 credit for each of dir1, dir2, dir3 containing file1."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('dir1_has_file1'):
        score += 1 / 3
    if result.get('dir2_has_file1'):
        score += 1 / 3
    if result.get('dir3_has_file1'):
        score += 1 / 3
    return min(round(score, 6), 1.0)

def check_timezone_newyork__cc3639ca98d2a9b7b2beb613c1282bfd_qw35sft2_6a4edbdc(result, expected, **options):
    """Check if system timezone is set to America/New_York (US Eastern Time).

    Parses `timedatectl status` output (line index 3 is the Time zone line).
    Example line: '                Time zone: America/New_York (EDT, -0400)'
    """
    if result is None:
        return 0.0
    try:
        lines = str(result).strip().split('\n')
        if len(lines) < 4:
            return 0.0
        tz_line = lines[3]
        colon_idx = tz_line.find(':')
        if colon_idx == -1:
            return 0.0
        tz_part = tz_line[colon_idx + 1:].strip()
        tz_name = tz_part.split(' ')[0]
        return 1.0 if 'New_York' in tz_name else 0.0
    except Exception:
        return 0.0

def check_new_files_perms__2e6f4a565bdf8ce8ae9bb4e56e1c79f7_qw35sft2_421bdfed(result, expected, **options):
    """Partial credit: 0.33 archive in old_files, 0.34 recent files moved to new_files, 0.33 new_files dir has expected permissions."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('archive_exists'):
        score += 0.33
    if result.get('new_files_populated'):
        score += 0.34
    expected_perm = expected.get('new_files_perm', '755')
    if result.get('new_files_perm') == expected_perm:
        score += 0.33
    return round(score, 2)

def check_home_users_file__baeeb37009528a1d856ed9e685fd5d89_qw35sft2_df5f2a52(result, expected, **options):
    """Check ~/home_users.txt exists and contains expected /etc/passwd entries."""
    if result.get('error') or not result.get('content'):
        return 0.0
    content = result['content']
    keywords = expected.get('keywords', [])
    if all((kw in content for kw in keywords)):
        return 1.0
    return 0.0

def check_text_scaling_exact__68667bdf1fe16ced5c6ec107d5f22e2d_qw35sft2_92c01ea7(result, expected, **options):
    """Check that text-scaling-factor matches expected_factor exactly (within 0.01 tolerance)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    try:
        factor = float(str(result).strip())
        target = expected.get('expected_factor', 1.25)
        return 1.0 if abs(factor - target) < 0.01 else 0.0
    except (ValueError, TypeError):
        return 0.0

def check_vim_removed_files_added__fafd25459d3a377f8d0c6d91b93474de_qw35sft2_5c623c89(result, expected, **options):
    """
    Partial-credit metric for: remove vim.desktop AND add org.gnome.Nautilus.desktop (Files).
    - 0.5 for vim.desktop absent from favorites
    - 0.5 for add_app (from expected rules) present in favorites
    Returns float in [0.0, 1.0].
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    apps = result.get('apps', [])
    if not isinstance(apps, list):
        return 0.0
    score = 0.0
    if 'vim.desktop' not in apps:
        score += 0.5
    add_app = expected.get('add_app', '')
    alt_app = expected.get('add_app_alt', '')
    if add_app and (add_app in apps or (alt_app and alt_app in apps)):
        score += 0.5
    return score

def check_rename_and_file__caa3a14501f9d1aa7da658fbd50adb72_qw35sft2_0b8869db(result, expected, **options):
    """
    Partial-credit check:
      0.5 - Desktop folder renamed to todo_list_Jan_2
      0.5 - notes.txt file created inside todo_list_Jan_2
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('renamed'):
        score += 0.5
    if result.get('file_exists'):
        score += 0.5
    return score

def check_file_perms_and_archive__00b9eddcbcb085dc89e5e17c3b839024_qw35sft2_ecf3946f(result, expected, **options):
    """Check files are 644 (0.5 pts) and testDir.tar archive exists (0.5 pts)."""
    if not result:
        return 0.0
    score = 0.0
    file_perms = result.get('file_permissions', [])
    expected_perm = expected.get('expected_file_permission', '644')
    if file_perms and all((p == expected_perm for p in file_perms)):
        score += 0.5
    if result.get('archive_exists', False):
        score += 0.5
    return score

def check_profile_renamed__c4720772b334cf5b9a4c5019ec6c389a_qw35sft2_d712d113(result, expected, **options):
    """Check that the 'default' profile was renamed in Thunderbird's profiles.ini."""
    if result is None:
        return 0.0
    content = str(result)
    new_name = expected.get('new_profile_name', 'backup')
    old_name = expected.get('old_profile_name', 'default')
    if not re.search(f'^Name={re.escape(new_name)}\\s*$', content, re.MULTILINE):
        return 0.0
    if re.search(f'^Name={re.escape(old_name)}\\s*$', content, re.MULTILINE):
        return 0.0
    return 1.0

def check_new_profile_created__0f2e8c08317dc15425e98e9b4b1171de_qw35sft2_aa3af12e(result, expected, **options):
    """Check that a new profile with the expected name was created in profiles.ini."""
    if result is None:
        return 0.0
    content = str(result)
    expected_name = expected.get('expected_profile_name', 'WorkProfile')
    if re.search(f'^Name={re.escape(expected_name)}\\s*$', content, re.MULTILINE):
        return 1.0
    return 0.0

def check_default_profile_changed__e7275df5a8446b39f1de20d975873334_qw35sft2_ac1506d0(result, expected, **options):
    """Check that the 'default' profile is marked as the default in profiles.ini.

    In Thunderbird's profiles.ini the default profile section contains 'Default=1'.
    We locate the section with 'Name=default' (exact match) and verify it has 'Default=1'.
    """
    if result is None:
        return 0.0
    content = str(result)
    target_name = expected.get('target_profile_name', 'default')
    sections = re.split('\\[Profile\\d+\\]', content)
    for section in sections:
        if re.search(f'^Name={re.escape(target_name)}\\s*$', section, re.MULTILINE):
            if re.search('^Default=1\\s*$', section, re.MULTILINE):
                return 1.0
            return 0.0
    return 0.0

def check_file_exists__df188ad33aeda3e315e62cc2c6afb173_qw35sft2_a5bd6840(result, expected, **options):
    """Return 1.0 if the target file exists in the VM."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('exists') else 0.0
