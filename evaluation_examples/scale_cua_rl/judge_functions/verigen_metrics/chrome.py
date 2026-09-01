"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_url_patterns__78f57355f208a4cab7262b3b71d760c1', 'check_bookmark_contains__4f204bdf47ab63e0c0d3c2d7f8cbe638', 'check_table_deleted__ffa5734cd742b977903cbddf387ddf6c', 'check_search_engine__29a0b35353e6c09dfbb5b1666d07a250', 'check_url_patterns__c6c4ea24b100fc1984ee26ee1f762d67', 'check_amazon_search_and_sort__68d5458a8c27a6ef63b63368a040c1bb', 'check_bookmark_and_active_tab__a0bd39d93c2a735cb19f1b8174d16ac1', 'check_timetable_entry__9319f3dc4336ee9ea3553360571f0046', 'check_url_patterns__99816a825a9ee8b512133b0316114c79', 'check_timetable_edit__929c2cac5efeaad2cf5e1556b5f1f38f', 'check_bookmark_folder_with_urls__1eba9f123a42adb66832822e9cf027cd', 'check_url_filters__c943b11f811330084211b7bb4efc17c3', 'check_url_contains__a5eff305d3fedc2c22bac974fb8698d7', 'check_chrome_startup_setting__8434e9ea78e4611ce770a316aad99fa9', 'check_search_engine__3898b17c9447f0df83eb7e4316fa07a5', 'check_bookmark_folders_subset__50cc0236f8d2558ccc3fa0de1fd3e9f1', 'check_html_file_keywords__0312c898008bbe83c35bf8d3f0838a0a', 'check_tabs_navigation__84216ceb2a30068e766e68b7d025a88e', 'check_html_title__43b4c11e2d677db2c4e2a819d325a693', 'check_hotel_search_sort__472d035cab9ac094961bab26b2de0108', 'check_extension_installed__c324ab1edeaa19edae46a63b788543c8', 'check_table_dimensions__126b20975e4e6ed136cbb81488fd0810', 'check_search_engine__5f4fc73f7b634bc96eb191d9ae7f5689', 'check_table_rows__d7904b6a5ae76e528f93447180e8419d', 'check_chrome_url__9f2f3fb3a73693af5f879fd7c403da86', 'check_table_row__d2dd546208d48aeec32a49334e180bbc', 'check_timetable_entry__ae3105aee7e1d6920d1ca91156ffa145', 'check_table_bg__9da3068d76a130c43b0820c0580d6baa', 'check_two_united_tabs__251fb4ef36f54773f246bbd000a7e5ba_qw35sft2_889ac4e7', 'check_chrome_delete_and_safebrowsing__be8c0e2643b53348f7e3a53d09b2fd8d_qw35sft2_c1bdde3a', 'check_chrome_dnt_and_lang__6c4389c0ed5d5685554268779bf7b6e2_qw35sft2_391472bb', 'check_chrome_profile_bookmark__f302c5cc6d17b170f6df5e52c54310d0_qw35sft2_71416571', 'check_tripadvisor_bookmarked__bf8a75ebb082889412df47eab97023bf_qw35sft2_f850de54', 'check_extensions_dev_mode__b759dda996122d64d46867fb7e10f84a_qw35sft2_cfe80f1e', 'check_third_party_cookies_blocked__2e738e6837968a7fb46f2d87cbf75561_qw35sft2_b88f5024', 'check_search_engine_bookmarks__7963cff65d97676543d55b7ec9b158c6_qw35sft2_c14ccc91', 'check_open_tabs_contain__13da459809674057ec3d167f8eb7ebb9_qw35sft2_2fd13702', 'check_new_tab_opened__8917e1d4c51d011a392d93603d714d31_qw35sft2_d8b52c64', 'check_chrome_profile_search__bae1be09edcf3b3e43698d058f3c0784_qw35sft2_7b076f1c', 'check_macys_bookmarked__9f1ee853f4a3b1e9c4f75b164c08128d_qw35sft2_928aea7a', 'check_url_contains__0e5c12133f110cb0d0d1de4db53f68cf_qw35sft2_1ae02646', 'check_chrome_delete_and_startup__f2a35c2c45e7134d422e18016f6e0ff7_qw35sft2_5900c37d', 'check_unpacked_extension_loaded__83ea9c43b50b1e067c01ffbea0552893_qw35sft2_f7b00670', 'check_google_search_toggle__d5be297d19eda82733e6c2f634aba9e2_qw35sft2_711e8ea3', 'check_bookmark_folder_and_dnt__33cb989983633f67656c84eb11865091_qw35sft2_0fe829cc', 'check_search_engine_dnt__67085c82c40b9ce759bdc0050139fed9_qw35sft2_c2c52651', 'check_url_contains__85e6b5742502db8c11d294296390de8f_qw35sft2_d7194790', 'check_multi_bookmark__34222c255c19ca63c5f6efa3ee3ba731_qw35sft2_c7fe46ee', 'check_chrome_delete_and_dnt__147c27e083ffc1ed15c90d12c99f1fef_qw35sft2_1fcac7f2', 'check_darktable_installed__bbd0b6cc65b6ed22ef194ee8993ef563_qw35sft2_9c2242fe', 'check_table_with_text__e5c61d88a7cca7248552b922fe7c3602_qw35sft2_0fc4d88d', 'check_table_dimensions__99738e8866681b9fdc97e255c87e124d_qw35sft2_4fac6981', 'check_two_tabstops__5d3adaf927fd3613be8530bc92e14de1_qw35sft2_3fa5d743', 'check_planet_table_rows__5135726bf1592f829eb43617d4e867fa_qw35sft2_1d980b93', 'check_right_tab_at_pos__14c984754ad8a9510e10b62f00a4c316_qw35sft2_8b02fccb', 'check_right_tab_at_pos__319c26794056ac5ca6a78fe65ba6a022_qw35sft2_1dfa36f5', 'check_table_7x5_inserted__94b525d2310c47dd52401ee93f279d0a_qw35sft2_d0c62e9f', 'check_left_tab_at_pos__cf559754b29c29c8081fb54264f1ac10_qw35sft2_bba96f0b', 'check_table_5x4_inserted__f2fed6259df9a3cde848f19734b96eb2_qw35sft2_7cd9effb', 'check_imdb_top250_bookmarked__27137be4f649298821e5e5376fdad889_qw35sft2_305df6b1', 'check_bookmark_url__05ce943d8d8f8a475d42f3e635259257_qw35sft2_5347a102', 'check_timetable_remove_tutorial6__c0f66209d7bea00756990ac3a15deab7_qw35sft2_647147b1', 'check_gpt4_table__8ae851f2fbec3e62c3a4b8c28bfb6c8b_qw35sft2_8ed65a33', 'check_link_opened_and_bookmarked__ba148d88a9035c4354cf728e6417abab_qw35sft2_6f6deeb5', 'check_github_python_tabs__a39b940fe4bed5c035657eb8fdc52896_qw35sft2_5b403118', 'check_timetable_lec_color__462e6cf4d0037b83d0687915bbffba64_qw35sft2_dbfe0525', 'check_table_and_methods__5d4910d5e528b018db8af7469dad985a_qw35sft2_4193fe88', 'check_extension_project_state__ce07ffbd7d3847464df6caec281dec8d_qw35sft2_ea0ad83c', 'check_vim_removed_chrome_first__41851786639af4a56b165470f26fecab_qw35sft2_62dc418d', 'check_troubleshooting_tab__ed5329c18a777573256ccdeb0cddcd91_qw35sft2_f17188d1', 'check_about_profiles_tab__8400d93b759737a4516cb53a7dfbb29f_qw35sft2_37b022de', 'check_both_extensions__b44386cdc8847efed8921bc79e3068ee_qw35sft2_bfdde70f']

def check_url_patterns__78f57355f208a4cab7262b3b71d760c1(result, expected, **options):
    """Check URL contains expected patterns with partial credit.

    Each pattern matched contributes equally to the score.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    url = result if isinstance(result, str) else str(result)
    url_lower = url.lower()
    patterns = expected.get('patterns', [])
    if not patterns:
        return 0.0
    score = 0.0
    per_pattern = 1.0 / len(patterns)
    for pattern in patterns:
        if re.search(pattern.lower(), url_lower):
            score += per_pattern
    return min(score, 1.0)

def check_bookmark_contains__4f204bdf47ab63e0c0d3c2d7f8cbe638(result, expected, **options):
    """Check if bookmarks contain a URL with the expected substring."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    url_substring = expected.get('url_substring', '')
    if not url_substring:
        return 0.0
    stack = [result]
    while stack:
        node = stack.pop()
        if isinstance(node, dict):
            if 'url' in node and url_substring in node.get('url', ''):
                return 1.0
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return 0.0

def check_table_deleted__ffa5734cd742b977903cbddf387ddf6c(result, expected, **options):
    """Check that the table has been deleted from the slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    should_exist = expected.get('table_should_exist', False)
    actual_exists = result.get('table_exists', True)
    if actual_exists == should_exist:
        return 1.0
    return 0.0

def check_search_engine__29a0b35353e6c09dfbb5b1666d07a250(result, expected, **options):
    """Check if the default search engine matches expected value."""
    if not result or not expected:
        return 0.0
    expected_engine = expected.get('expected_engine', '')
    if isinstance(result, str) and result.strip().lower() == expected_engine.strip().lower():
        return 1.0
    return 0.0

def check_url_patterns__c6c4ea24b100fc1984ee26ee1f762d67(result, expected, **options):
    """Check URL contains expected patterns with partial credit.

    Each pattern matched contributes equally to the score.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    url = result if isinstance(result, str) else str(result)
    url_lower = url.lower()
    patterns = expected.get('patterns', [])
    if not patterns:
        return 0.0
    score = 0.0
    per_pattern = 1.0 / len(patterns)
    for pattern in patterns:
        if re.search(pattern.lower(), url_lower):
            score += per_pattern
    return min(score, 1.0)

def check_amazon_search_and_sort__68d5458a8c27a6ef63b63368a040c1bb(result, expected, **options):
    """Check if the browser navigated to Amazon search results with correct query and sort order.

    Partial credit:
    - 0.5 for correct search term in URL
    - 0.5 for correct sort parameter in URL
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    url = result if isinstance(result, str) else str(result)
    score = 0.0
    search_pattern = expected.get('search_pattern', '')
    if search_pattern and re.search(search_pattern, url, re.IGNORECASE):
        score += 0.5
    sort_pattern = expected.get('sort_pattern', '')
    if sort_pattern and re.search(sort_pattern, url, re.IGNORECASE):
        score += 0.5
    return min(score, 1.0)

def check_bookmark_and_active_tab__a0bd39d93c2a735cb19f1b8174d16ac1(result, expected, **options):
    """Check bookmark exists in bar AND active tab matches expected URL.

    Partial credit:
    - 0.5 for bookmark present in bookmarks bar
    - 0.5 for correct active tab

    Args:
        result: dict with 'bookmark_urls' (list) and 'active_url' (str)
        expected: dict with 'bookmark_url' (str) and 'active_tab_url' (str)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_bookmark = expected.get('bookmark_url', '')
    bookmark_urls = result.get('bookmark_urls', [])
    for url in bookmark_urls:
        if expected_bookmark and expected_bookmark in url:
            score += 0.5
            break
    expected_active = expected.get('active_tab_url', '')
    actual_active = result.get('active_url', '')
    if expected_active and expected_active in actual_active:
        score += 0.5
    return min(score, 1.0)

def check_timetable_entry__9319f3dc4336ee9ea3553360571f0046(result, expected, **options):
    """Check if a timetable cell has the expected text keywords and background color.

    Partial credit:
    - 0.5 for correct text (contains required keywords)
    - 0.5 for correct background color
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    value = result.get('value')
    bg_color = result.get('bg_color', 'none')
    required_keywords = expected.get('required_keywords', [])
    if value and required_keywords:
        value_str = str(value).lower()
        all_found = all((kw.lower() in value_str for kw in required_keywords))
        if all_found:
            score += 0.5
    expected_color = expected.get('expected_color', '')
    if expected_color and bg_color != 'none':
        if expected_color.upper() in bg_color.upper():
            score += 0.5
    return min(score, 1.0)

def check_url_patterns__99816a825a9ee8b512133b0316114c79(result, expected, **options):
    """Check URL contains expected patterns with partial credit.

    Each pattern matched contributes equally to the score.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    url = result if isinstance(result, str) else str(result)
    url_lower = url.lower()
    patterns = expected.get('patterns', [])
    if not patterns:
        return 0.0
    score = 0.0
    per_pattern = 1.0 / len(patterns)
    for pattern in patterns:
        if re.search(pattern.lower(), url_lower):
            score += per_pattern
    return min(score, 1.0)

def check_timetable_edit__929c2cac5efeaad2cf5e1556b5f1f38f(result, expected, **options):
    """Check if a timetable cell was renamed and recolored correctly.

    Partial credit:
    - 0.5 for correct text (contains required keywords, does NOT contain forbidden keywords)
    - 0.5 for correct background color
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    value = result.get('value')
    bg_color = result.get('bg_color', 'none')
    required_keywords = expected.get('required_keywords', [])
    forbidden_keywords = expected.get('forbidden_keywords', [])
    if value:
        value_str = str(value).lower()
        has_required = all((kw.lower() in value_str for kw in required_keywords)) if required_keywords else False
        has_forbidden = any((kw.lower() in value_str for kw in forbidden_keywords)) if forbidden_keywords else False
        if has_required and (not has_forbidden):
            score += 0.5
    expected_color = expected.get('expected_color', '')
    if expected_color and bg_color != 'none':
        if expected_color.upper() in bg_color.upper():
            score += 0.5
    return min(score, 1.0)

def check_bookmark_folder_with_urls__1eba9f123a42adb66832822e9cf027cd(result, expected, **options):
    """
    Check bookmark folder exists and contains expected URL.
    Partial credit:
      - 0.5 for folder existing
      - 0.5 for correct URL in folder
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('folder_exists', False):
        score += 0.5
        logger.info('Folder exists: +0.5')
        expected_url_substring = expected.get('expected_url_substring', '')
        folder_urls = result.get('folder_urls', [])
        for url in folder_urls:
            if expected_url_substring and expected_url_substring in url:
                score += 0.5
                logger.info(f"Found URL containing '{expected_url_substring}': +0.5")
                break
    else:
        logger.info('Folder does not exist')
    return min(score, 1.0)

def check_url_filters__c943b11f811330084211b7bb4efc17c3(result, expected, **options):
    """Check if URL contains expected filter patterns with partial credit.

    Scoring:
    - 0.5 for short sleeve filter applied
    - 0.5 for size L filter applied
    Prerequisite: must be on men's shirts page (otherwise 0.0).
    """
    if not result or (isinstance(result, str) and 'error' in result.lower()):
        return 0.0
    url = result if isinstance(result, str) else str(result)
    score = 0.0
    if not re.search('macys\\.com.*mens.*shirts', url, re.IGNORECASE):
        return 0.0
    if re.search('Short.*Sleeve', url, re.IGNORECASE):
        score += 0.5
    if re.search('Men_regular_size_t', url) and re.search('[/,]L[,?/&]|[/,]L$', url):
        score += 0.5
    return min(score, 1.0)

def check_url_contains__a5eff305d3fedc2c22bac974fb8698d7(result, expected, **options):
    """Check if the browser URL contains the expected substring."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_url = result.get('url', '')
    if not actual_url:
        return 0.0
    required_parts = expected.get('required_url_parts', [])
    if not required_parts:
        return 0.0
    matched = 0
    for part in required_parts:
        if part.lower() in actual_url.lower():
            matched += 1
    return matched / len(required_parts)

def check_chrome_startup_setting__8434e9ea78e4611ce770a316aad99fa9(result, expected, **options):
    """Check if Chrome startup setting matches expected restore_on_startup value.

    Scoring:
    - 1.0 if restore_on_startup matches expected value
    - 0.0 otherwise
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_restore = expected.get('expected_restore_on_startup')
    actual_restore = result.get('restore_on_startup')
    if actual_restore == expected_restore:
        return 1.0
    return 0.0

def check_search_engine__3898b17c9447f0df83eb7e4316fa07a5(result, expected, **options):
    """Check if default search engine matches expected value.

    Scoring:
    - 1.0 if search engine name matches (case-insensitive)
    - 0.0 otherwise
    """
    if not result or (isinstance(result, dict) and result.get('error')):
        return 0.0
    expected_engine = expected.get('expected_engine', '')
    if isinstance(result, str):
        if result.lower().strip() == expected_engine.lower().strip():
            return 1.0
    return 0.0

def check_bookmark_folders_subset__50cc0236f8d2558ccc3fa0de1fd3e9f1(result, expected, **options):
    """Check that expected folder names are a subset of actual bookmark bar folder names.

    Unlike the official is_expected_bookmarks which uses exact set equality,
    this metric uses subset checking to handle pre-existing folders on the bookmarks bar.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    rule = expected
    if rule.get('type') != 'bookmark_bar_folders_names':
        return 0.0
    try:
        bookmark_bar = result.get('bookmark_bar', {})
        children = bookmark_bar.get('children', [])
        actual_folder_names = set((child['name'] for child in children if child.get('type') == 'folder'))
        expected_names = set(rule.get('names', []))
        if expected_names.issubset(actual_folder_names):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_html_file_keywords__0312c898008bbe83c35bf8d3f0838a0a(result, expected, **options):
    """Check if HTML file exists, has sufficient size, and contains expected keywords.

    Scoring:
    - 0.3: File exists
    - 0.2: File exceeds minimum size
    - 0.5: Keywords found in content (proportional)
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    score += 0.3
    min_size = expected.get('min_size', 100)
    if result.get('size', 0) > min_size:
        score += 0.2
    content = result.get('content_lower', '')
    keywords = expected.get('keywords', [])
    if keywords:
        found = sum((1 for kw in keywords if kw.lower() in content))
        keyword_score = found / len(keywords) * 0.5
        score += keyword_score
    return min(score, 1.0)

def check_tabs_navigation__84216ceb2a30068e766e68b7d025a88e(result, expected, **options):
    """Check if open tabs contain expected URL substrings. Partial credit per match."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    tab_urls = []
    if isinstance(result, list):
        for tab in result:
            if isinstance(tab, dict):
                tab_urls.append(tab.get('url', ''))
            elif isinstance(tab, str):
                tab_urls.append(tab)
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    score = 0.0
    points_per_check = 1.0 / len(checks)
    for check in checks:
        substring = check.get('url_contains', '')
        weight = check.get('weight', points_per_check)
        if substring and any((substring in url for url in tab_urls)):
            score += weight
    return min(score, 1.0)

def check_html_title__43b4c11e2d677db2c4e2a819d325a693(result, expected, **options):
    """Check HTML file exists and has correct title. Partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title', '')
    if expected_title and expected_title.lower() == actual_title.lower():
        score += 0.5
    return min(score, 1.0)

def check_hotel_search_sort__472d035cab9ac094961bab26b2de0108(result, expected, **options):
    """Check NYC hotel search with sort by price. Partial credit: 0.5 for city, 0.5 for sort."""
    score = 0.0
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    actual_city = str(result.get('city', ''))
    expected_city = expected.get('expected_city', '')
    if expected_city and expected_city in actual_city:
        score += 0.5
    actual_rank = str(result.get('rank', ''))
    expected_rank = expected.get('expected_rank', '')
    if expected_rank and expected_rank in actual_rank:
        score += 0.5
    return min(score, 1.0)

def check_extension_installed__c324ab1edeaa19edae46a63b788543c8(result, expected, **options):
    """Check if an expected VS Code extension is installed."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    extensions = result.get('extensions', '')
    expected_ext = expected.get('expected_extension', '')
    if expected_ext.lower() in extensions.lower():
        return 1.0
    return 0.0

def check_table_dimensions__126b20975e4e6ed136cbb81488fd0810(result, expected, **options):
    """Check if a new table was inserted with expected dimensions.

    Scoring:
    - 0.5 for having a new table (table_count increased)
    - 0.25 for correct number of columns
    - 0.25 for correct number of rows
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_table_count', 2)
    expected_rows = expected.get('expected_rows', 6)
    expected_cols = expected.get('expected_cols', 4)
    actual_count = result.get('table_count', 0)
    actual_rows = result.get('last_table_rows', 0)
    actual_cols = result.get('last_table_cols', 0)
    if actual_count >= expected_count:
        score += 0.5
    if actual_cols == expected_cols:
        score += 0.25
    if actual_rows == expected_rows:
        score += 0.25
    return min(score, 1.0)

def check_search_engine__5f4fc73f7b634bc96eb191d9ae7f5689(result, expected, **options):
    """Check if Chrome's default search engine matches expected value.

    Args:
        result: String value from default_search_engine getter (e.g., "Google", "DuckDuckGo").
        expected: Dict with 'expected_engine' key.

    Returns:
        float: 1.0 if result matches expected engine (case-insensitive), 0.0 otherwise.
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    expected_engine = expected.get('expected_engine', '')
    if isinstance(result, str) and isinstance(expected_engine, str):
        if result.strip().lower() == expected_engine.strip().lower():
            return 1.0
    return 0.0

def check_table_rows__d7904b6a5ae76e528f93447180e8419d(result, expected, **options):
    """Check if the table has the expected number of rows and the last row has correct values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('expected_row_count', 0)
    actual_rows = result.get('row_count', 0)
    row_count_ok = actual_rows >= expected_rows
    if row_count_ok:
        score += 0.5
    expected_last_row = expected.get('expected_last_row', [])
    actual_last_row = result.get('last_row_values', [])
    if expected_last_row and actual_last_row:
        if [v.strip() for v in actual_last_row] == [v.strip() for v in expected_last_row]:
            score += 0.5
    return min(score, 1.0)

def check_chrome_url__9f2f3fb3a73693af5f879fd7c403da86(result, expected, **options):
    """Check if Chrome active URL contains expected domain."""
    if not result or isinstance(result, str) or result.get('error'):
        return 0.0
    actual_url = result.get('url', '')
    expected_url = expected.get('expected_url', '')
    if not actual_url or not expected_url:
        return 0.0
    if expected_url in actual_url:
        return 1.0
    return 0.0

def check_table_row__d2dd546208d48aeec32a49334e180bbc(result, expected, **options):
    """Check if a table row matches expected values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    total = len(expected_values)
    if total == 0:
        return 0.0
    matches = 0
    for (i, exp_val) in enumerate(expected_values):
        if i < len(actual_values):
            if str(actual_values[i]).strip() == str(exp_val).strip():
                matches += 1
    return matches / total

def check_timetable_entry__ae3105aee7e1d6920d1ca91156ffa145(result, expected, **options):
    """Check if a timetable cell has the expected text keywords and background color.

    Partial credit:
    - 0.5 for correct text (contains required keywords)
    - 0.5 for correct background color
    """
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    value = result.get('value')
    bg_color = result.get('bg_color', 'none')
    required_keywords = expected.get('required_keywords', [])
    if value and required_keywords:
        value_str = str(value).lower()
        all_found = all((kw.lower() in value_str for kw in required_keywords))
        if all_found:
            score += 0.5
    expected_color = expected.get('expected_color', '')
    if expected_color and bg_color != 'none':
        if expected_color.upper() in bg_color.upper():
            score += 0.5
    return min(score, 1.0)

def check_table_bg__9da3068d76a130c43b0820c0580d6baa(result, expected, **options):
    """Check table cell value and background color. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_value = expected.get('expected_cell_value', '')
    actual_value = result.get('cell_value', '')
    if expected_value and actual_value:
        if expected_value.lower() == actual_value.lower():
            score += 0.5
    expected_colors = expected.get('expected_bg_colors', [])
    actual_color = result.get('bg_color', '')
    if actual_color and expected_colors:
        actual_upper = actual_color.upper()
        for ec in expected_colors:
            if ec.upper() == actual_upper:
                score += 0.5
                break
    return min(score, 1.0)

def check_two_united_tabs__251fb4ef36f54773f246bbd000a7e5ba_qw35sft2_889ac4e7(result, expected, **options):
    """Check that two specified United Airlines pages are open as separate tabs.

    Partial credit: 0.5 per matching tab (total 1.0 if both present).

    Args:
        result: List of tab info dicts with 'url' key (from get_open_tabs_info getter).
        expected: Rules dict with 'tab1_pattern' and 'tab2_pattern' substrings.

    Returns:
        float in [0.0, 1.0].
    """
    if not result or isinstance(result, str):
        return 0.0
    tab1_pattern = expected.get('tab1_pattern', '')
    tab2_pattern = expected.get('tab2_pattern', '')
    if not tab1_pattern or not tab2_pattern:
        return 0.0
    urls = []
    for tab in result:
        if isinstance(tab, dict):
            urls.append(tab.get('url', ''))
        elif isinstance(tab, (list, tuple)) and len(tab) >= 1:
            urls.append(str(tab[0]))
    score = 0.0
    if any((tab1_pattern in url for url in urls)):
        score += 0.5
    if any((tab2_pattern in url for url in urls)):
        score += 0.5
    return score

def check_chrome_delete_and_safebrowsing__be8c0e2643b53348f7e3a53d09b2fd8d_qw35sft2_c1bdde3a(result, expected, **options):
    """Partial credit: 0.5 for auto-delete on close, 0.5 for Enhanced Safe Browsing enabled."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('data_delete') == expected.get('data_delete'):
        score += 0.5
    if result.get('enhanced_safe_browsing') == expected.get('enhanced_safe_browsing'):
        score += 0.5
    return score

def check_chrome_dnt_and_lang__6c4389c0ed5d5685554268779bf7b6e2_qw35sft2_391472bb(result, expected, **options):
    """Check DNT is enabled (0.5) and language matches expected prefix (0.5)."""
    if result is None or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    if result.get('dnt') == expected.get('dnt', True):
        score += 0.5
    expected_lang = expected.get('language', '')
    actual_lang = str(result.get('language', '')).lower()
    if expected_lang and actual_lang.startswith(expected_lang.lower()):
        score += 0.5
    return score

def check_chrome_profile_bookmark__f302c5cc6d17b170f6df5e52c54310d0_qw35sft2_71416571(result, expected, **options):
    """Check Chrome profile name and presence of a bookmark URL. Partial credit: 0.5 each."""
    if result is None or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    if result.get('profile_name') == expected.get('profile_name'):
        score += 0.5
    expected_host = expected.get('bookmark_host', '')
    bookmark_bar_urls = result.get('bookmark_bar_urls', [])
    if expected_host and any((expected_host in url for url in bookmark_bar_urls)):
        score += 0.5
    return score

def check_tripadvisor_bookmarked__bf8a75ebb082889412df47eab97023bf_qw35sft2_f850de54(result, expected, **options):
    """Check if TripAdvisor is saved as a bookmark in Chrome."""
    if not result:
        return 0.0
    expected_substr = expected.get('expected_url_substr', 'tripadvisor.com')
    stack = [result]
    while stack:
        node = stack.pop()
        if isinstance(node, dict):
            url = node.get('url', '')
            if url and expected_substr in url:
                return 1.0
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return 0.0

def check_extensions_dev_mode__b759dda996122d64d46867fb7e10f84a_qw35sft2_cfe80f1e(result, expected, **options):
    """Check if Chrome Developer mode matches the expected boolean value."""
    if result is None or isinstance(result, str):
        return 0.0
    expected_val = expected.get('expected', True)
    actual_val = result.get('developer_mode', False)
    return 1.0 if bool(actual_val) == bool(expected_val) else 0.0

def check_third_party_cookies_blocked__2e738e6837968a7fb46f2d87cbf75561_qw35sft2_b88f5024(result, expected, **options):
    """Check if Chrome is set to block all third-party cookies (mode == 2)."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    mode = result.get('cookie_controls_mode', -1)
    if mode == 2:
        return 1.0
    return 0.0

def check_search_engine_bookmarks__7963cff65d97676543d55b7ec9b158c6_qw35sft2_c14ccc91(result, expected, **options):
    """Check Bing is default search engine (0.5) and bing.com is in the bookmarks bar (0.5). Partial credit."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_engine = expected.get('expected_engine', 'Bing')
    if expected_engine.lower() in str(result.get('search_engine', '')).lower():
        score += 0.5
    if result.get('bing_bookmarked') == expected.get('expected_bookmarked', True):
        score += 0.5
    return score

def check_open_tabs_contain__13da459809674057ec3d167f8eb7ebb9_qw35sft2_2fd13702(result, expected, **options):
    """Check that a tab with the expected domain is present among open tabs."""
    if not isinstance(result, list):
        return 0.0
    expected_domain = expected.get('expected_domain', '')
    if not expected_domain:
        return 0.0
    for tab in result:
        url = tab.get('url', '') or ''
        if expected_domain in url:
            return 1.0
    return 0.0

def check_new_tab_opened__8917e1d4c51d011a392d93603d714d31_qw35sft2_d8b52c64(result, expected, **options):
    """Check that at least min_tabs tabs are open in Chrome.

    result: list of {title, url} dicts from open_tabs_info getter.
    expected: dict with 'min_tabs' key (default 2).
    """
    if not isinstance(result, list):
        return 0.0
    min_tabs = expected.get('min_tabs', 2)
    return 1.0 if len(result) >= min_tabs else 0.0

def check_chrome_profile_search__bae1be09edcf3b3e43698d058f3c0784_qw35sft2_7b076f1c(result, expected, **options):
    """Check Chrome profile name and default search engine. Partial credit: 0.5 each."""
    if result is None or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    if result.get('profile_name') == expected.get('profile_name'):
        score += 0.5
    actual_engine = (result.get('search_engine') or '').strip()
    expected_engine = (expected.get('search_engine') or '').strip()
    if actual_engine.lower() == expected_engine.lower():
        score += 0.5
    return score

def check_macys_bookmarked__9f1ee853f4a3b1e9c4f75b164c08128d_qw35sft2_928aea7a(result, expected, **options):
    """Check if macys.com is bookmarked in the bookmarks bar.

    result: dict returned by get_bookmarks (nested bookmarks structure)
    expected: dict from rules, e.g. {"target_url": "macys.com"}
    """
    target = expected.get('target_url', 'macys.com')
    if not result or not isinstance(result, dict):
        return 0.0
    bookmark_bar = result.get('bookmark_bar', result.get('Bookmarks bar', {}))
    stack = [bookmark_bar]
    while stack:
        node = stack.pop()
        if isinstance(node, dict):
            url = node.get('url', '')
            if url and target in url:
                return 1.0
            for val in node.values():
                if isinstance(val, (dict, list)):
                    stack.append(val)
        elif isinstance(node, list):
            for item in node:
                stack.append(item)
    return 0.0

def check_url_contains__0e5c12133f110cb0d0d1de4db53f68cf_qw35sft2_1ae02646(result, expected, **options):
    """Check if the current URL matches all required patterns (variation 0: Hotels section)."""
    if result is None:
        return 0.0
    url = result if isinstance(result, str) else str(result)
    patterns = expected.get('url_patterns', [])
    if not patterns:
        return 0.0
    for pattern in patterns:
        if not re.search(pattern, url, re.IGNORECASE):
            return 0.0
    return 1.0

def check_chrome_delete_and_startup__f2a35c2c45e7134d422e18016f6e0ff7_qw35sft2_5900c37d(result, expected, **options):
    """Partial credit: 0.5 for auto-delete on close, 0.5 for continue-where-left-off startup."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('data_delete') == expected.get('data_delete'):
        score += 0.5
    if result.get('restore_on_startup') == expected.get('restore_on_startup'):
        score += 0.5
    return score

def check_unpacked_extension_loaded__83ea9c43b50b1e067c01ffbea0552893_qw35sft2_f7b00670(result, expected, **options):
    """Check if a specific unpacked extension is loaded in Chrome.

    If extension_name is provided, checks for a case-insensitive substring match.
    If extension_name is empty, checks that at least one unpacked extension is loaded.
    """
    if result is None or isinstance(result, str):
        return 0.0
    expected_name = expected.get('extension_name', '')
    actual_names = result.get('unpacked_extensions', [])
    if not expected_name:
        return 1.0 if len(actual_names) > 0 else 0.0
    for name in actual_names:
        if expected_name.lower() in str(name).lower():
            return 1.0
    return 0.0

def check_google_search_toggle__d5be297d19eda82733e6c2f634aba9e2_qw35sft2_711e8ea3(result, expected, **options):
    """Check if a Google Search settings toggle is in the expected state.

    expected keys:
        expected_state (bool): True = ON, False = OFF
    Returns 1.0 if match, 0.0 otherwise.
    """
    if result is None:
        return 0.0
    if isinstance(result, dict) and result.get('error') and (result.get('checked') is None):
        return 0.0
    actual = None
    if isinstance(result, dict):
        actual = result.get('checked')
    if actual is None:
        return 0.0
    expected_state = expected.get('expected_state', True)
    return 1.0 if bool(actual) == bool(expected_state) else 0.0

def check_bookmark_folder_and_dnt__33cb989983633f67656c84eb11865091_qw35sft2_0fe829cc(result, expected, **options):
    """Check that a named bookmark folder exists and Do Not Track is enabled.

    Partial credit: 0.5 for folder present, 0.5 for Do Not Track enabled.
    expected keys: folder_name (str), do_not_track (bool)
    """
    score = 0.0
    folder_name = expected.get('folder_name', 'Favorites')
    expected_dnt = expected.get('do_not_track', True)
    if folder_name in result.get('folder_names', []):
        score += 0.5
    if result.get('do_not_track', False) == expected_dnt:
        score += 0.5
    return score

def check_search_engine_dnt__67085c82c40b9ce759bdc0050139fed9_qw35sft2_c2c52651(result, expected, **options):
    """Check Bing is default search engine (0.5) and Do Not Track is enabled (0.5). Partial credit."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    score = 0.0
    expected_engine = expected.get('expected_engine', 'Bing')
    if expected_engine.lower() in str(result.get('search_engine', '')).lower():
        score += 0.5
    if result.get('do_not_track') == expected.get('expected_dnt', 'true'):
        score += 0.5
    return score

def check_url_contains__85e6b5742502db8c11d294296390de8f_qw35sft2_d7194790(result, expected, **options):
    """Check if the current URL matches all required patterns (variation 2: Restaurants section)."""
    if result is None:
        return 0.0
    url = result if isinstance(result, str) else str(result)
    patterns = expected.get('url_patterns', [])
    if not patterns:
        return 0.0
    for pattern in patterns:
        if not re.search(pattern, url, re.IGNORECASE):
            return 0.0
    return 1.0

def check_multi_bookmark__34222c255c19ca63c5f6efa3ee3ba731_qw35sft2_c7fe46ee(result, expected, **options):
    """Check that multiple URLs are present in the bookmarks bar. Partial credit per URL."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    bar_urls = result.get('bar_urls', [])
    required_urls = expected.get('required_urls', [])
    if not required_urls:
        return 0.0
    credit_per_url = 1.0 / len(required_urls)
    score = 0.0
    for url in required_urls:
        if any((url in bar_url or bar_url in url for bar_url in bar_urls)):
            score += credit_per_url
    return min(round(score, 4), 1.0)

def check_chrome_delete_and_dnt__147c27e083ffc1ed15c90d12c99f1fef_qw35sft2_1fcac7f2(result, expected, **options):
    """Partial credit: 0.5 for auto-delete on close, 0.5 for Do Not Track enabled."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('data_delete') == expected.get('data_delete'):
        score += 0.5
    if result.get('do_not_track') == expected.get('do_not_track'):
        score += 0.5
    return score

def check_darktable_installed__bbd0b6cc65b6ed22ef194ee8993ef563_qw35sft2_9c2242fe(result, expected, **options):
    """Check that darktable is installed on the system."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('installed'):
        return 1.0
    return 0.0

def check_table_with_text__e5c61d88a7cca7248552b922fe7c3602_qw35sft2_0fc4d88d(result, expected, **options):
    """Check that a table with correct dimensions exists and first cell has expected text."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_rows = expected.get('rows', 5)
    expected_cols = expected.get('cols', 2)
    expected_text = expected.get('first_cell_text', '')
    tables = result.get('tables', [])
    if not tables:
        return 0.0
    score = 0.0
    for tbl in tables:
        if tbl.get('rows') == expected_rows and tbl.get('cols') == expected_cols:
            score += 0.7
            if expected_text:
                actual_text = tbl.get('first_cell_text', '')
                if expected_text.lower() in actual_text.lower():
                    score += 0.3
            else:
                score += 0.3
            return min(score, 1.0)
    return 0.0

def check_table_dimensions__99738e8866681b9fdc97e255c87e124d_qw35sft2_4fac6981(result, expected, **options):
    """Check that a table with specified rows and cols exists on the slide."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_rows = expected.get('rows', 6)
    expected_cols = expected.get('cols', 3)
    tables = result.get('tables', [])
    if not tables:
        return 0.0
    for tbl in tables:
        if tbl.get('rows') == expected_rows and tbl.get('cols') == expected_cols:
            return 1.0
    return 0.0

def check_two_tabstops__5d3adaf927fd3613be8530bc92e14de1_qw35sft2_3fa5d743(result, expected, **options):
    """Check that all non-empty paragraphs have both a left tab at ~0 cm and a right tab at ~15 cm.
    Partial credit: 0.5 per requirement (left tab presence + right tab presence), averaged over paragraphs."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    left_pos = float(expected.get('left_position_cm', 0.0))
    right_pos = float(expected.get('right_position_cm', 15.0))
    tolerance = float(expected.get('tolerance_cm', 0.5))
    left_aligns = {'LEFT', 'START'}
    right_aligns = {'RIGHT', 'END'}
    paragraphs = result.get('paragraphs', [])
    if not paragraphs:
        return 0.0
    total_score = 0.0
    for para in paragraphs:
        stops = para.get('stops', [])
        has_left = any((abs(s['position_cm'] - left_pos) <= tolerance and s['alignment'].upper() in left_aligns for s in stops))
        has_right = any((abs(s['position_cm'] - right_pos) <= tolerance and s['alignment'].upper() in right_aligns for s in stops))
        total_score += 0.5 * int(has_left) + 0.5 * int(has_right)
    return total_score / len(paragraphs)

def check_planet_table_rows__5135726bf1592f829eb43617d4e867fa_qw35sft2_1d980b93(result, expected, **options):
    """Check whether the planet comparison table contains a Pluto row with expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_pluto', False):
        score += 0.5
    pluto_row = result.get('pluto_row')
    expected_diameter = expected.get('expected_diameter', '2,376')
    if pluto_row and len(pluto_row) > 1:
        diameter_cell = pluto_row[1]
        if expected_diameter.replace(',', '') in diameter_cell.replace(',', ''):
            score += 0.25
    expected_mass = expected.get('expected_mass', '0.013')
    if pluto_row and len(pluto_row) > 2:
        mass_cell = pluto_row[2]
        if expected_mass in mass_cell:
            score += 0.25
    return min(score, 1.0)

def check_right_tab_at_pos__14c984754ad8a9510e10b62f00a4c316_qw35sft2_8b02fccb(result, expected, **options):
    """Check that all non-empty paragraphs have a right-aligned tab stop at the expected position."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required_pos = float(expected.get('position_cm', 15.0))
    required_align = expected.get('alignment', 'RIGHT').upper()
    tolerance = float(expected.get('tolerance_cm', 0.5))
    align_aliases = {'RIGHT': {'RIGHT', 'END'}, 'LEFT': {'LEFT', 'START'}}
    valid_aligns = align_aliases.get(required_align, {required_align})
    paragraphs = result.get('paragraphs', [])
    if not paragraphs:
        return 0.0
    matching = 0
    for para in paragraphs:
        for stop in para.get('stops', []):
            pos_ok = abs(stop['position_cm'] - required_pos) <= tolerance
            align_ok = stop['alignment'].upper() in valid_aligns
            if pos_ok and align_ok:
                matching += 1
                break
    return matching / len(paragraphs)

def check_right_tab_at_pos__319c26794056ac5ca6a78fe65ba6a022_qw35sft2_1dfa36f5(result, expected, **options):
    """Check that all non-empty paragraphs have a right-aligned tab stop at the expected position."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required_pos = float(expected.get('position_cm', 10.0))
    required_align = expected.get('alignment', 'RIGHT').upper()
    tolerance = float(expected.get('tolerance_cm', 0.5))
    align_aliases = {'RIGHT': {'RIGHT', 'END'}, 'LEFT': {'LEFT', 'START'}}
    valid_aligns = align_aliases.get(required_align, {required_align})
    paragraphs = result.get('paragraphs', [])
    if not paragraphs:
        return 0.0
    matching = 0
    for para in paragraphs:
        for stop in para.get('stops', []):
            pos_ok = abs(stop['position_cm'] - required_pos) <= tolerance
            align_ok = stop['alignment'].upper() in valid_aligns
            if pos_ok and align_ok:
                matching += 1
                break
    return matching / len(paragraphs)

def check_table_7x5_inserted__94b525d2310c47dd52401ee93f279d0a_qw35sft2_d0c62e9f(result, expected, **options):
    """Check that a new 7-column x 5-row table was inserted (document now has 2 tables).
    Partial credit: 0.5 for correct table count, 0.5 for correct dimensions of new table.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count') == expected.get('table_count', 2):
        score += 0.5
    if result.get('last_rows') == expected.get('last_rows', 5) and result.get('last_cols') == expected.get('last_cols', 7):
        score += 0.5
    return score

def check_left_tab_at_pos__cf559754b29c29c8081fb54264f1ac10_qw35sft2_bba96f0b(result, expected, **options):
    """Check that all non-empty paragraphs have a left-aligned tab stop at the expected position."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required_pos = float(expected.get('position_cm', 5.0))
    required_align = expected.get('alignment', 'LEFT').upper()
    tolerance = float(expected.get('tolerance_cm', 0.5))
    align_aliases = {'RIGHT': {'RIGHT', 'END'}, 'LEFT': {'LEFT', 'START'}}
    valid_aligns = align_aliases.get(required_align, {required_align})
    paragraphs = result.get('paragraphs', [])
    if not paragraphs:
        return 0.0
    matching = 0
    for para in paragraphs:
        for stop in para.get('stops', []):
            pos_ok = abs(stop['position_cm'] - required_pos) <= tolerance
            align_ok = stop['alignment'].upper() in valid_aligns
            if pos_ok and align_ok:
                matching += 1
                break
    return matching / len(paragraphs)

def check_table_5x4_inserted__f2fed6259df9a3cde848f19734b96eb2_qw35sft2_7cd9effb(result, expected, **options):
    """Check that a new 5-column x 4-row table was inserted (document now has 2 tables).
    Partial credit: 0.5 for correct table count, 0.5 for correct dimensions.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count') == expected.get('table_count', 2):
        score += 0.5
    if result.get('last_rows') == expected.get('last_rows', 4) and result.get('last_cols') == expected.get('last_cols', 5):
        score += 0.5
    return score

def check_imdb_top250_bookmarked__27137be4f649298821e5e5376fdad889_qw35sft2_305df6b1(result, expected, **options):
    """Check if any bookmark contains the IMDB Top 250 chart URL."""
    if result is None:
        return 0.0
    expected_url_fragment = expected.get('expected_url_fragment', 'imdb.com/chart/top')
    stack = [result]
    while stack:
        node = stack.pop()
        if isinstance(node, dict):
            url = node.get('url', '')
            if url and expected_url_fragment in url:
                return 1.0
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return 0.0

def check_bookmark_url__05ce943d8d8f8a475d42f3e635259257_qw35sft2_5347a102(result, expected, **options):
    """Check that a specific URL is saved in Chrome bookmarks."""
    if not result:
        return 0.0
    target_url = expected.get('expected_url', '')
    if not target_url:
        return 0.0
    stack = [result]
    while stack:
        node = stack.pop()
        if isinstance(node, dict):
            url = node.get('url', '')
            if url and target_url in url:
                return 1.0
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return 0.0

def check_timetable_remove_tutorial6__c0f66209d7bea00756990ac3a15deab7_qw35sft2_647147b1(result, expected, **options):
    """Check D5 has Wed 12PM lecture (0.5) and G12 is empty/removed (0.5)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    d5_value = result.get('d5_value') or ''
    if '12:00' in d5_value or 'Lec 2' in d5_value or 'lec 2' in d5_value.lower():
        score += 0.5
    g12_value = result.get('g12_value')
    if g12_value is None or g12_value == 'None' or g12_value.strip() == '':
        score += 0.5
    return score

def check_gpt4_table__8ae851f2fbec3e62c3a4b8c28bfb6c8b_qw35sft2_8ed65a33(result, expected, **options):
    """Return 1.0 if a table containing the GPT-4 row exists in the '5.2 Main Results' section."""
    import re
    if result.get('error') or not result.get('has_table'):
        return 0.0
    model_name = expected.get('model_name', 'Gpt-4')
    elements = result.get('elements', [])
    if not elements:
        for table in result.get('tables', []):
            for row in table:
                if any((model_name.lower() in str(cell).lower() for cell in row)):
                    return 1.0
        return 0.0
    section_idx = None
    for i, elem in enumerate(elements):
        if elem.get('type') == 'paragraph':
            text = elem.get('text', '')
            if '5.2' in text and 'main results' in text.lower():
                section_idx = i
                break
    if section_idx is None:
        for i, elem in enumerate(elements):
            if elem.get('type') == 'paragraph':
                text = elem.get('text', '')
                if 'main results' in text.lower():
                    section_idx = i
                    break
    if section_idx is None:
        return 0.0
    section_heading_text = elements[section_idx].get('text', '')
    next_section_idx = len(elements)
    for i in range(section_idx + 1, len(elements)):
        elem = elements[i]
        if elem.get('type') == 'paragraph':
            text = elem.get('text', '')
            if text and re.match('^\\d+(\\.\\d+)*\\s+\\w', text) and (text != section_heading_text):
                next_section_idx = i
                break
    for i in range(section_idx + 1, next_section_idx):
        elem = elements[i]
        if elem.get('type') == 'table':
            table_data = elem.get('data', [])
            for row in table_data:
                if any((model_name.lower() in str(cell).lower() for cell in row)):
                    return 1.0
    return 0.0

def check_link_opened_and_bookmarked__ba148d88a9035c4354cf728e6417abab_qw35sft2_6f6deeb5(result, expected, **options):
    """
    Partial-credit metric: checks both that a URL was opened in Chrome (any tab)
    AND that it was saved to Chrome bookmarks.
    Score: 0.5 for tab presence + 0.5 for bookmark presence.
    """
    if not result or result.get('error'):
        return 0.0
    target_url = expected.get('expected_url', '')
    if not target_url:
        return 0.0
    score = 0.0
    for tab in result.get('tabs', []):
        if target_url in tab.get('url', ''):
            score += 0.5
            break
    stack = [result.get('bookmarks', {})]
    found = False
    while stack and (not found):
        node = stack.pop()
        if isinstance(node, dict):
            if target_url in node.get('url', ''):
                found = True
            else:
                stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    if found:
        score += 0.5
    return score

def check_github_python_tabs__a39b940fe4bed5c035657eb8fdc52896_qw35sft2_5b403118(result, expected, **options):
    """Check if both github.com and docs.python.org are open in Chrome tabs.

    result: list of {title, url} dicts from open_tabs_info getter
    expected: dict with github_url and python_url keys (already unwrapped from rules)
    Returns: 0.5 per found URL, max 1.0
    """
    if not result or not isinstance(result, list):
        return 0.0
    github_host = expected.get('github_url', 'github.com')
    python_host = expected.get('python_url', 'docs.python.org')
    all_urls = [tab.get('url', '') for tab in result if isinstance(tab, dict)]
    github_found = any((github_host in url for url in all_urls))
    python_found = any((python_host in url for url in all_urls))
    score = 0.0
    if github_found:
        score += 0.5
    if python_found:
        score += 0.5
    return score

def check_timetable_lec_color__462e6cf4d0037b83d0687915bbffba64_qw35sft2_dbfe0525(result, expected, **options):
    """Check D5 has a Wed 12PM lecture entry (0.5) and yellow background (0.5)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    d5_value = result.get('d5_value') or ''
    if '12:00' in d5_value or 'Lec 2' in d5_value or 'lec 2' in d5_value.lower():
        score += 0.5
    d5_color = result.get('d5_fill_color') or ''
    expected_color = expected.get('expected_fill_color', 'FFFFFF00')
    if d5_color == expected_color:
        score += 0.5
    return score

def check_table_and_methods__5d4910d5e528b018db8af7469dad985a_qw35sft2_4193fe88(result, expected, **options):
    """Partial credit: 0.5 for GPT-4 table, 0.5 for GPT-4 mention in 5.1 Methods section."""
    if result.get('error'):
        return 0.0
    score = 0.0
    model_name = expected.get('model_name', 'Gpt-4')
    for table in result.get('tables', []):
        for row in table:
            if any((model_name.lower() in str(cell).lower() for cell in row)):
                score += 0.5
                break
        if score >= 0.5:
            break
    required_text = expected.get('methods_text', 'gpt-4')
    paragraphs = result.get('paragraphs', [])
    in_section = False
    for para in paragraphs:
        if '5.1' in para and 'Methods' in para:
            in_section = True
            continue
        if in_section:
            if para and para.startswith('5.2'):
                break
            if para and para.lower() != 'todo' and (required_text.lower() in para.lower()):
                score += 0.5
                break
    return min(score, 1.0)

def check_extension_project_state__ce07ffbd7d3847464df6caec281dec8d_qw35sft2_ea0ad83c(result, expected, **options):
    """Partial-credit check for happy-extension project structure in ~/Projects.

    Scoring:
      - 0.34: ~/Projects/happy-extension/ directory exists
      - 0.33: manifest.json exists inside the directory
      - 0.33: background_script.js exists inside the directory
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('dir'):
        score += 0.34
    if result.get('manifest'):
        score += 0.33
    if result.get('background'):
        score += 0.33
    return min(score, 1.0)

def check_vim_removed_chrome_first__41851786639af4a56b165470f26fecab_qw35sft2_62dc418d(result, expected, **options):
    """
    Partial-credit metric for: remove vim.desktop AND make google-chrome.desktop first in favorites.
    - 0.5 for vim.desktop absent from favorites
    - 0.5 for first_app (from expected rules) being the first entry in favorites
    Returns float in [0.0, 1.0].
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    apps = result.get('apps', [])
    if not isinstance(apps, list):
        return 0.0
    score = 0.0
    if 'vim.desktop' not in apps:
        score += 0.5
    first_app = expected.get('first_app', '')
    if first_app and apps and (apps[0] == first_app):
        score += 0.5
    return score

def check_troubleshooting_tab__ed5329c18a777573256ccdeb0cddcd91_qw35sft2_f17188d1(result, expected, **options):
    """Check if the Troubleshooting Information page is open in Thunderbird."""
    if result is None:
        return 0.0
    tree_str = str(result)
    expected_text = expected.get('expected_text', 'Troubleshooting Information')
    return 1.0 if expected_text in tree_str else 0.0

def check_about_profiles_tab__8400d93b759737a4516cb53a7dfbb29f_qw35sft2_37b022de(result, expected, **options):
    """Check if the About Profiles page is open as an active tab in Thunderbird."""
    if result is None:
        return 0.0
    tree_str = str(result)
    expected_text = expected.get('expected_text', 'About Profiles')
    return 1.0 if expected_text in tree_str else 0.0

def check_both_extensions__b44386cdc8847efed8921bc79e3068ee_qw35sft2_bfdde70f(result, expected, **options):
    """Check both Python and Jupyter extensions are installed. 0.5 credit per extension."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if isinstance(result, dict):
        ext_str = result.get('output', '') or result.get('stdout', '') or str(result)
    else:
        ext_str = str(result) if result else ''
    python_ext = expected.get('python_ext', 'ms-python.python')
    jupyter_ext = expected.get('jupyter_ext', 'ms-toolsai.jupyter')
    score = 0.0
    if python_ext in ext_str:
        score += 0.5
    if jupyter_ext in ext_str:
        score += 0.5
    return score
