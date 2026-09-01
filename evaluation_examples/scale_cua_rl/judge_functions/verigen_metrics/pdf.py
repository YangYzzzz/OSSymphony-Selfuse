"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_pdf_exists__18b0663c1e76805fd4ea486cbf72380b', 'check_pdf_list__f72fa74134cf5a58d1281947ac1e62bd', 'check_lecture_pdf_exists__29ea2d20dd9405b64f3df8f5049b79d3', 'check_ods_to_pdf__40ad56f2b3098a6972e48b06b58c95e0', 'check_pdf_export__beb4b25746cd97a98dc95327476aaf4d', 'check_pdf_valid__1c15db0bd188a0b79dfb455b9076c68e', 'check_pdf_count__358385204fa7b84b2d9118b3314be461', 'check_pdf_orientation__075106d955ead60b31af19408ecd54ac', 'check_downloads_pdf__5416755d2dd223fd56c07a64ea26e507_qw35sft2_5ed1afcc', 'check_csv_and_pdf__f45249b0314207a95d7366d15bee2907_qw35sft2_c790450a', 'check_pdf_exists__19cbc30a1547517beef14a189cce767e_qw35sft2_69d39035', 'check_orgsummary_pdf_exists__82b855263e1bdeb6cb9c2f640e53bf6c_qw35sft2_3689ff13', 'check_pdf_in_documents__df1f48650ab6f100a1208b8040cd8828_qw35sft2_ca603cf2', 'check_chapter_pdf_moved__6964660b6570aa960a367d963aeb2477_qw35sft2_ff5568f0', 'check_lecture_pdf_downloaded__ac6b45f9bcb788e83171ced5593fc0f6_qw35sft2_89c3dc97', 'check_lecture_pdf_downloaded__6f1f6989222847d5b9f522129093f325_qw35sft2_ed434618']

def check_pdf_exists__18b0663c1e76805fd4ea486cbf72380b(result, expected, **options):
    """Check if the PDF file exists and has a minimum size."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    min_size = expected.get('min_size', 1000)
    if result.get('size', 0) >= min_size:
        return 1.0
    return 0.0

def check_pdf_list__f72fa74134cf5a58d1281947ac1e62bd(result, expected, **options):
    """Check if PDF list file contains all expected filenames. Partial credit per file."""
    if not isinstance(result, dict) or result.get('error') or (not result.get('exists')):
        return 0.0
    content_lower = result.get('content_lower', '')
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    found = 0
    for fname in expected_files:
        if fname.lower() in content_lower:
            found += 1
    return found / len(expected_files)

def check_lecture_pdf_exists__29ea2d20dd9405b64f3df8f5049b79d3(result, expected, **options):
    """Check if lecture PDF exists and has reasonable size."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False):
        min_size = expected.get('min_size', 1000)
        if result.get('size', 0) >= min_size:
            return 1.0
        return 0.5
    return 0.0

def check_ods_to_pdf__40ad56f2b3098a6972e48b06b58c95e0(result, expected, **options):
    """Check ODS to PDF conversion: terminal usage + valid PDF output."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('used_terminal') == 'use terminal':
        score += 0.5
    min_size = expected.get('expected_min_size', 1000)
    if result.get('file_exists') and result.get('is_valid_pdf'):
        if result.get('file_size', 0) >= min_size:
            score += 0.5
        else:
            score += 0.25
    return min(score, 1.0)

def check_pdf_export__beb4b25746cd97a98dc95327476aaf4d(result, expected, **options):
    """Check if PDF was exported correctly."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error') or not result.get('exists'):
        return 0.0
    score = 0.0
    if result.get('is_pdf'):
        score += 0.5
    expected_pages = expected.get('expected_pages', 17)
    actual_pages = result.get('page_count', 0)
    if actual_pages == expected_pages:
        score += 0.5
    elif actual_pages > 0:
        score += 0.2
    return min(score, 1.0)

def check_pdf_valid__1c15db0bd188a0b79dfb455b9076c68e(result, expected, **options):
    """Check if PDF file exists and is valid."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.5
    if result.get('is_pdf'):
        score += 0.3
    if result.get('size', 0) > 1024:
        score += 0.2
    return min(score, 1.0)

def check_pdf_count__358385204fa7b84b2d9118b3314be461(result, expected, **options):
    """Check if text file contains the expected count number."""
    if result.get('error'):
        return 0.0
    content = result.get('content', '').strip()
    expected_count = str(expected.get('expected_count', ''))
    import re
    numbers = re.findall('\\d+', content)
    if expected_count in numbers:
        return 1.0
    if content == expected_count:
        return 1.0
    return 0.0

def check_pdf_orientation__075106d955ead60b31af19408ecd54ac(result, expected, **options):
    """Check if PDF exists, has expected orientation, and has no margins."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('file_exists'):
        score += 0.34
    else:
        return 0.0
    expected_landscape = expected.get('is_landscape', True)
    if result.get('is_landscape') == expected_landscape:
        score += 0.33
    expected_no_margins = expected.get('has_no_margins', True)
    if result.get('margin_check_available', False):
        if result.get('has_no_margins') == expected_no_margins:
            score += 0.33
    else:
        pass
    return min(score, 1.0)

def check_downloads_pdf__5416755d2dd223fd56c07a64ea26e507_qw35sft2_5ed1afcc(result, expected, **options):
    """Check PDF in Downloads with correct name (0.5) and page count for no-margins (0.5).

    Trajectory evidence: saving with margins=None produces 22 pages vs 23 for default margins.
    This page-count proxy is used to verify the no-margin requirement.
    """
    if not result or not isinstance(result, dict):
        return 0.0
    pdf_files = result.get('pdf_files', [])
    if not pdf_files:
        return 0.0
    pdf_name_contains = expected.get('pdf_name_contains', '')
    if pdf_name_contains:
        found = any((pdf_name_contains.lower() in f.lower() for f in pdf_files))
        if not found:
            return 0.0
    score = 0.5
    expected_page_count = expected.get('expected_page_count')
    if expected_page_count is not None:
        actual_page_count = result.get('page_count')
        if actual_page_count is not None and actual_page_count == expected_page_count:
            score += 0.5
    else:
        score += 0.5
    return score

def check_csv_and_pdf__f45249b0314207a95d7366d15bee2907_qw35sft2_c790450a(result, expected, **options):
    """Partial credit: 0.5 for CSV exported, 0.5 for PDF exported."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('csv_exists') and result.get('csv_row_count', 0) >= 4:
        score += 0.5
    if result.get('pdf_exists') and result.get('pdf_size', 0) > 100:
        score += 0.5
    return score

def check_pdf_exists__19cbc30a1547517beef14a189cce767e_qw35sft2_69d39035(result, expected, **options):
    """Return 1.0 if the PDF file exists and has content, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False) and result.get('size', 0) > 0:
        return 1.0
    return 0.0

def check_orgsummary_pdf_exists__82b855263e1bdeb6cb9c2f640e53bf6c_qw35sft2_3689ff13(result, expected, **options):
    """Return 1.0 if OrgSummary.pdf exists with content, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False) and result.get('size', 0) > 0:
        return 1.0
    return 0.0

def check_pdf_in_documents__df1f48650ab6f100a1208b8040cd8828_qw35sft2_ca603cf2(result, expected, **options):
    """Return 1.0 if PDF exists in Documents folder with content, 0.0 otherwise."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False) and result.get('size', 0) > 0:
        return 1.0
    return 0.0

def check_chapter_pdf_moved__6964660b6570aa960a367d963aeb2477_qw35sft2_ff5568f0(result, expected, **options):
    """Check that the chapter 1 PDF was moved from book folder to Desktop.

    Gives full credit if file is present at destination and absent at original.
    Gives partial credit (0.5) if file exists at destination regardless of original.
    """
    if not isinstance(result, dict):
        return 0.0
    file_at_dest = result.get('file_at_dest', False)
    file_at_orig = result.get('file_at_orig', True)
    if file_at_dest and (not file_at_orig):
        return 1.0
    if file_at_dest:
        return 0.5
    return 0.0

def check_lecture_pdf_downloaded__ac6b45f9bcb788e83171ced5593fc0f6_qw35sft2_89c3dc97(result, expected, **options):
    """Check if the required PDF file was downloaded to lecture_slides."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    files = result.get('files', []) if isinstance(result, dict) else []
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    found = sum((1 for f in required_files if f in files))
    return found / len(required_files)

def check_lecture_pdf_downloaded__6f1f6989222847d5b9f522129093f325_qw35sft2_ed434618(result, expected, **options):
    """Check if the required PDF file was downloaded to lecture_slides."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    files = result.get('files', []) if isinstance(result, dict) else []
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    found = sum((1 for f in required_files if f in files))
    return found / len(required_files)
