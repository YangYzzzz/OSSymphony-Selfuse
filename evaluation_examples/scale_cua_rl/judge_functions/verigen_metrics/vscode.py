"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_workspace_folders__56abb1eed13e878c335e4f175376e020', 'check_vscode_setting__ce8e3f26fe80873188ba40e45449d3b9', 'check_vscode_setting__215b1cc2c3da9fb70c240f9f4f8df6b5', 'check_workspace_setting__50a1f916fbc8eceedebd80aa7a09697c', 'check_vscode_setting__b082db988cea152fff3c88fce0623d97', 'check_vscode_settings__24dcbea28334e3c805df8afddc88a5f1', 'check_workspace_setting__6a6725db0c0fb336a2ee8d61d5b62877', 'check_vscode_setting__c7225cce9efdba9aa8af211dee6285ef', 'check_vscode_settings_multi__ebe9d7bea41ec6a21d7329851f18b701', 'check_vscode_setting__ef57df2dacd824ebd85cb86fda295eab', 'check_vscode_workspace_open__2e0260d8e24ca4d49ed462b4c3422ce2', 'check_full_workspace_bookmark__218f219ceca7e7e0003dc76221985da3_qw35sft2_26aeae6a', 'check_vscode_debug_focus_console__b4f4f2d972bad7dfc4a5dbad64bc0fb1_qw35sft2_3bde8aeb', 'check_vscode_keybinding__ed42af8e24d428ca40b99f135c68274f_qw35sft2_bb493910', 'check_vscode_multi_settings__81417cbab160b1f4e52b4e40c1e83805_qw35sft2_c3a35ae8', 'check_vscode_setting__e20226d1d49056af54e3e69873c8fc3c_qw35sft2_38ff2d0b', 'check_vscode_dual_settings__95e499e1a1bb4bed6d6d6a23e9ae7a30_qw35sft2_a25bcb21', 'check_vscode_settings__193e70599e6295621968803241fc32b6_qw35sft2_9f9b5635', 'check_vscode_word_wrap__15748b95140e1ac4229209dbc7dfe067_qw35sft2_1082a02e', 'check_vscode_remove_two_list_shortcuts__45d3929cc132a565457f8b673ec620af_qw35sft2_d0c81f2e', 'check_workspace_and_file__5eb515badd2fc6ca56c5f84fbe245437_qw35sft2_bae3d1a2', 'check_vscode_wrap_and_tab__5464614ee51132440430b0764a874f9a_qw35sft2_e78900d6', 'check_vscode_exclude_dual__26878716b564c56320e14f3b9550ebe4_qw35sft2_a6e31fe0', 'check_vscode_dual_ext__eb3392bef87727a9af1f8c1971adb226_qw35sft2_6440db73', 'check_vscode_file_and_settings__9fe429b95a40aa1badcb7de0b8a6b0f1_qw35sft2_87823a8f', 'check_vscode_locale_and_tabsize__de00e991b11dacec5fa1917014fc1ed1_qw35sft2_b04b42e3', 'check_vscode_py_indent_docstring__3d465870a880c6e96493a5637b57636e_qw35sft2_9804e750', 'check_vscode_multi_settings__8bfa6a8d7119af313a0660a3dea4ceab_qw35sft2_bc343dc4', 'check_vscode_keybinding__70885a60e67c4def20a9d0e3284c1793_qw35sft2_500cb562', 'check_vscode_dual_settings__c64cfd038604e0824dbc8809a35f407d_qw35sft2_830dc09c', 'check_vscode_setting__fa93c32e8afd83b84bab5f585818b0fb_qw35sft2_c69a4670', 'check_ext_and_keybinding__44f31fc07357666acac95e9a419ee84f_qw35sft2_dbb9caad', 'check_vscode_theme_wordwrap__978e7cba9602912bc979b36fdc172116_qw35sft2_b6f8aa43', 'check_workspace_has_folders__f72e9e4f450564b5101c9194862d205b_qw35sft2_548bf84b', 'check_vscode_settings__d45b70bac3694dd736206b10dff26100_qw35sft2_ed066c92', 'check_vscode_exclude_pytest__6a66711a30fb269a0e7d77ca8b497959_qw35sft2_fd948bea', 'check_vscode_remove_and_add_keybinding__8bf5980383e9d80a451eac2b89ccf1df_qw35sft2_30e84c9c', 'check_vscode_word_wrap_settings__35c11665b6e5c9d9c9abe7f855ccfe46_qw35sft2_91d69975', 'check_vscode_ext_autosave__e2317900e4eb6fe8b0284d167aa38270_qw35sft2_3c369c77', 'check_workspace_and_extension__bead7b293cda299d38413ae3ff873cc8_qw35sft2_1a5e842f', 'check_vscode_py_function_file__9e0f9c32752df4ad0956b05cabebbc98_qw35sft2_e4880cfc', 'check_vscode_workspace_in_storage__5752e6d2cd4f13c8d24f076655028c23_qw35sft2_5ec3117b', 'check_vscode_py_indent_blankline__2c836d0eca9976b79042c49c51fe71cf_qw35sft2_62042e83', 'check_vscode_multi_settings__9ec2182c98e7d1d45bbbb3ae2d0db42e_qw35sft2_feab433d', 'check_vscode_debug_focus_wrap__09cc1b24020e8f7c0f7b48c6a4d96413_qw35sft2_19312b4c', 'check_vscode_setting__a4b8eefcd32d2ac2b259c28c71814550_qw35sft2_8583a75e', 'check_workspace_has_folder__d5491cdef8176b887dce317a57139150_qw35sft2_0784154b', 'check_vscode_setting__5cd158cc213aefb25242b372a3862ff3_qw35sft2_e8a19b8e', 'check_vscode_wrap_and_format__95a613ea8825af09f40057c978f02ad7_qw35sft2_54a5ef43', 'check_vscode_locale__215b1cc2c3da9fb70c240f9f4f8df6b5_qw35sft2_c98e144e', 'check_workspace_and_wordwrap__1713b5c17b096bf267a520ed66d388f0_qw35sft2_8d1703c8', 'check_vscode_exclude_autosave__7797789b3c2086d41b9d5f1699195599_qw35sft2_b6097013', 'check_vscode_desktop_file_comment__eea98fddd6e0176c917193edbee2b203_qw35sft2_d3906e8a', 'check_vscode_dual_settings__37e5b967fa33d7c5b2d11883c3b2a4f6_qw35sft2_4811c951', 'check_vscode_ext_wordwrap__88ac2a64a7ee030ff707715de339ea34_qw35sft2_5cfa3a00', 'check_vscode_multi_settings__d122d575a22853e47d56ee92b6c18378_qw35sft2_33e4705c', 'check_vscode_open_workspace__c37846152daf06e05c18277c5b1f636d_qw35sft2_caafa045', 'check_workspace_folders_replaced__edd7ecbec0169642eef724f574a58652_qw35sft2_3273d2e6', 'vscode_workspace_saved__f91feeeb0a889e0b6b274960d3d06c7d_qw35sft2_280f3ee8', 'check_vscode_tab_size__015bdcf4a8f96215b5579fc53db2d526_qw35sft2_451e5295', 'check_vscode_debug_focus_terminal__29baecbd3036f18eff818341b1c0e4e8_qw35sft2_814b6d3c', 'check_vscode_setting__d97bffb72b4b8c27def4b087597b01a0_qw35sft2_f7177a16', 'check_vscode_dual_settings__b75b59f2415937a4444d85150c1ab46c_qw35sft2_5a8c7a44', 'check_vscode_single_ext__521eb992bf97de2d79e4aa6c081db684_qw35sft2_ba37d45a', 'check_vscode_active_file__6a582969bc6d0b9624887969279510e3_qw35sft2_33b62156', 'check_vscode_multi_settings__7e25498cca269e458737cc32ebf8beb0_qw35sft2_9903d5f1']

def check_workspace_folders__56abb1eed13e878c335e4f175376e020(result, expected, **options):
    """Check if workspace contains all expected folders. Supports partial credit."""
    if result.get('error'):
        return 0.0
    actual_folders = result.get('folders', [])
    expected_folders = expected.get('expected_folders', [])
    if not expected_folders:
        return 0.0
    score = 0.0
    per_folder = 1.0 / len(expected_folders)
    for ef in expected_folders:
        for af in actual_folders:
            if ef in af or af in ef:
                score += per_folder
                break
    return min(score, 1.0)

def check_vscode_setting__ce8e3f26fe80873188ba40e45449d3b9(result, expected, **options):
    """Check if VS Code settings match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected_settings', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, exp_val) in expected_settings.items():
        actual_val = settings.get(key)
        if actual_val == exp_val:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_setting__215b1cc2c3da9fb70c240f9f4f8df6b5(result, expected, **options):
    """Check if VS Code settings match expected values. Partial credit for multiple settings."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, value) in expected_settings.items():
        if settings.get(key) == value:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_workspace_setting__50a1f916fbc8eceedebd80aa7a09697c(result, expected, **options):
    """Check that workspace settings contain all expected key-value pairs (subset check)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected', {})
    if not expected_settings:
        return 0.0
    for (key, value) in expected_settings.items():
        if key not in settings or settings[key] != value:
            return 0.0
    return 1.0

def check_vscode_setting__b082db988cea152fff3c88fce0623d97(result, expected, **options):
    """Check if VS Code settings match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected_settings', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, exp_val) in expected_settings.items():
        actual_val = settings.get(key)
        if actual_val == exp_val:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_settings__24dcbea28334e3c805df8afddc88a5f1(result, expected, **options):
    """Check VS Code settings match expected values. Partial credit per setting."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected_settings', {})
    if not expected_settings:
        return 0.0
    total = len(expected_settings)
    matched = 0
    for (key, exp_value) in expected_settings.items():
        actual_value = settings.get(key)
        if actual_value == exp_value:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_workspace_setting__6a6725db0c0fb336a2ee8d61d5b62877(result, expected, **options):
    """Check if a specific workspace setting has the expected value."""
    if result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    setting_key = expected.get('setting_key', '')
    expected_value = expected.get('expected_value')
    actual_value = settings.get(setting_key)
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_vscode_setting__c7225cce9efdba9aa8af211dee6285ef(result, expected, **options):
    """Check if VS Code settings match expected values. Partial credit for multiple settings."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, value) in expected_settings.items():
        if settings.get(key) == value:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_settings_multi__ebe9d7bea41ec6a21d7329851f18b701(result, expected, **options):
    """Check if multiple VS Code settings match expected values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected_settings', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, exp_val) in expected_settings.items():
        actual_val = settings.get(key)
        if actual_val == exp_val:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_setting__ef57df2dacd824ebd85cb86fda295eab(result, expected, **options):
    """Check if VS Code settings match expected values. Partial credit for multiple settings."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_settings = expected.get('expected', {})
    if not expected_settings:
        return 0.0
    matched = 0
    total = len(expected_settings)
    for (key, value) in expected_settings.items():
        if settings.get(key) == value:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_workspace_open__2e0260d8e24ca4d49ed462b4c3422ce2(result, expected, **options):
    """Check if expected workspace name appears in VS Code window titles."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    titles = result.get('window_titles', '')
    keyword = expected.get('keyword', '')
    if keyword and keyword.lower() in titles.lower():
        return 1.0
    return 0.0

def check_full_workspace_bookmark__218f219ceca7e7e0003dc76221985da3_qw35sft2_26aeae6a(result, expected, **options):
    """Check full workspace: terminal in OSWorld, github.com open, docs.python.org open, Python docs bookmarked.

    result: dict with 'terminal_cwd', 'open_urls', 'bookmarks' keys
    expected: dict with 'expected_cwd', 'github_url', 'python_url', 'bookmark_host' keys
    Returns: 0.25 per satisfied condition, max 1.0
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_cwd = expected.get('expected_cwd', '/home/user/Documents/Projects/OSWorld')
    github_url = expected.get('github_url', 'github.com')
    python_url = expected.get('python_url', 'docs.python.org')
    bookmark_host = expected.get('bookmark_host', 'docs.python.org')
    terminal_cwd = result.get('terminal_cwd', '')
    open_urls = result.get('open_urls', [])
    bookmarks = result.get('bookmarks', {})
    score = 0.0
    if expected_cwd in terminal_cwd or terminal_cwd.endswith('OSWorld'):
        score += 0.25
    if any((github_url in url for url in open_urls)):
        score += 0.25
    if any((python_url in url for url in open_urls)):
        score += 0.25
    _stack = [bookmarks]
    _found_bookmark = False
    while _stack and (not _found_bookmark):
        _node = _stack.pop()
        if isinstance(_node, dict):
            if bookmark_host in _node.get('url', ''):
                _found_bookmark = True
            else:
                _stack.extend(_node.values())
        elif isinstance(_node, list):
            _stack.extend(_node)
    if _found_bookmark:
        score += 0.25
    return min(score, 1.0)

def check_vscode_debug_focus_console__b4f4f2d972bad7dfc4a5dbad64bc0fb1_qw35sft2_3bde8aeb(result, expected, **options):
    """Check debug.focusEditorOnBreak is False and debug.internalConsoleOptions is correct. 0.5 credit each."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('focus_editor_on_break') == expected.get('focus_editor_on_break'):
        score += 0.5
    if result.get('internal_console_options') == expected.get('internal_console_options'):
        score += 0.5
    return min(score, 1.0)

def check_vscode_keybinding__ed42af8e24d428ca40b99f135c68274f_qw35sft2_bb493910(result, expected, **options):
    """Check if VS Code keybindings.json contains the expected keybinding entry.

    expected keys:
        key     (str): the keyboard shortcut, e.g. 'ctrl+shift+j'
        command (str): the VS Code command id
    """
    if result.get('error') or not result.get('entries'):
        return 0.0
    entries = result.get('entries', [])
    expected_key = expected.get('key', '').lower().replace(' ', '')
    expected_command = expected.get('command', '')
    for entry in entries:
        if not isinstance(entry, dict):
            continue
        entry_key = entry.get('key', '').lower().replace(' ', '')
        entry_command = entry.get('command', '')
        if entry_command.startswith('-'):
            continue
        if entry_key == expected_key and entry_command == expected_command:
            return 1.0
    return 0.0

def check_vscode_multi_settings__81417cbab160b1f4e52b4e40c1e83805_qw35sft2_c3a35ae8(result, expected, **options):
    """Check multiple VS Code settings with partial credit.

    expected keys:
      checks: list of {"key": str, "value": <expected_value>}
    Scoring: equal weight per check.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    weight = 1.0 / len(checks)
    score = 0.0
    for check in checks:
        key = check['key']
        exp_value = check['value']
        actual = result.get(key)
        if actual == exp_value:
            score += weight
        elif isinstance(exp_value, (int, float)) and isinstance(actual, (int, float)):
            if float(actual) == float(exp_value):
                score += weight
    return min(round(score, 4), 1.0)

def check_vscode_setting__e20226d1d49056af54e3e69873c8fc3c_qw35sft2_38ff2d0b(result, expected, **options):
    """Check that a specific VS Code setting matches the expected value."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    setting_key = expected.get('setting_key', '')
    expected_value = expected.get('expected_value')
    if not setting_key:
        return 0.0
    actual_value = result.get(setting_key)
    if isinstance(expected_value, (int, float)) and isinstance(actual_value, (int, float)):
        return 1.0 if actual_value == expected_value else 0.0
    return 1.0 if actual_value == expected_value else 0.0

def check_vscode_dual_settings__95e499e1a1bb4bed6d6d6a23e9ae7a30_qw35sft2_a25bcb21(result, expected, **options):
    """Check two VS Code settings with partial credit (0.5 each).

    result: local path to the downloaded settings.json file (from vm_file getter)
    expected: rules dict already unwrapped by get_rule()
      - setting1_key: str
      - setting1_val: any
      - setting2_key: str
      - setting2_val: any
    """
    import json
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    score = 0.0
    setting1_key = expected.get('setting1_key')
    setting1_val = expected.get('setting1_val')
    setting2_key = expected.get('setting2_key')
    setting2_val = expected.get('setting2_val')
    if setting1_key is not None and settings.get(setting1_key) == setting1_val:
        score += 0.5
    if setting2_key is not None and settings.get(setting2_key) == setting2_val:
        score += 0.5
    return score

def check_vscode_settings__193e70599e6295621968803241fc32b6_qw35sft2_9f9b5635(result, expected, **options):
    """Check that VS Code settings.json contains the expected key-value pairs."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    for key, expected_val in expected.items():
        actual_val = result.get(key)
        if actual_val != expected_val:
            return 0.0
    return 1.0

def check_vscode_word_wrap__15748b95140e1ac4229209dbc7dfe067_qw35sft2_1082a02e(result, expected, **options):
    """Check that editor.wordWrap matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else None
    expected_value = expected.get('expected_value')
    if expected_value is None:
        return 0.0
    return 1.0 if actual == expected_value else 0.0

def check_vscode_remove_two_list_shortcuts__45d3929cc132a565457f8b673ec620af_qw35sft2_d0c81f2e(result, expected, **options):
    """Check both ctrl+f (list.find) and Escape (list.closeFind) are removed. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('ctrlf_removed'):
        score += 0.5
    if result.get('escape_removed'):
        score += 0.5
    return score

def check_workspace_and_file__5eb515badd2fc6ca56c5f84fbe245437_qw35sft2_bae3d1a2(result, expected, **options):
    """Partial-credit check: 0.5 for folder in workspace + 0.5 for file existence."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    folder_names = result.get('folder_names', [])
    required_folder = expected.get('required_folder_name', 'data1')
    if required_folder in folder_names:
        score += 0.5
    if result.get('file_exists', False):
        score += 0.5
    return min(score, 1.0)

def check_vscode_wrap_and_tab__5464614ee51132440430b0764a874f9a_qw35sft2_e78900d6(result, expected, **options):
    """Check wordWrapColumn and tabSize. 0.5 per sub-goal."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_col = result.get('wordWrapColumn')
    exp_col = expected.get('wordWrapColumn')
    if exp_col is not None and actual_col == exp_col:
        score += 0.5
    actual_tab = result.get('tabSize')
    exp_tab = expected.get('tabSize')
    if exp_tab is not None and actual_tab == exp_tab:
        score += 0.5
    return score

def check_vscode_exclude_dual__26878716b564c56320e14f3b9550ebe4_qw35sft2_a6e31fe0(result, expected, **options):
    """Check files.exclude contains pattern1 (0.5) and pattern2 (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    files_exclude = result.get('files_exclude', {})
    pattern1 = expected.get('pattern1', '')
    pattern2 = expected.get('pattern2', '')
    if pattern1 and files_exclude.get(pattern1):
        score += 0.5
    if pattern2 and files_exclude.get(pattern2):
        score += 0.5
    return score

def check_vscode_dual_ext__eb3392bef87727a9af1f8c1971adb226_qw35sft2_6440db73(result, expected, **options):
    """Check if two VS Code extensions are installed (partial credit 0.5 each)."""
    if not isinstance(result, dict):
        return 0.0
    ext_list = result.get('extensions', '')
    score = 0.0
    ext1 = expected.get('ext1', '').lower()
    ext2 = expected.get('ext2', '').lower()
    if ext1 and ext1 in ext_list:
        score += 0.5
    if ext2 and ext2 in ext_list:
        score += 0.5
    return score

def check_vscode_file_and_settings__9fe429b95a40aa1badcb7de0b8a6b0f1_qw35sft2_87823a8f(result, expected, **options):
    """Partial credit: 0.5 for test.py existing, 0.5 for tab size matching expected."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_tab_size = expected.get('expected_tab_size', 2)
    if result.get('file_exists'):
        score += 0.5
    tab_size = result.get('tab_size')
    if tab_size is not None and int(tab_size) == int(expected_tab_size):
        score += 0.5
    return min(score, 1.0)

def check_vscode_locale_and_tabsize__de00e991b11dacec5fa1917014fc1ed1_qw35sft2_b04b42e3(result, expected, **options):
    """Check VS Code display language is Chinese (Simplified) and editor.tabSize is 4.

    result: dict with 'locale' and 'settings' from getter
    expected: rules dict (already unwrapped), contains 'locale' and 'tab_size'
    Partial credit: 0.5 per sub-goal.
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_locale = expected.get('locale', '').lower()
    actual_locale = result.get('locale', '').lower()
    if expected_locale and (actual_locale == expected_locale or actual_locale.startswith(expected_locale)):
        score += 0.5
    expected_tab = expected.get('tab_size')
    if expected_tab is not None:
        settings = result.get('settings', {})
        actual_tab = settings.get('editor.tabSize')
        if actual_tab == expected_tab:
            score += 0.5
    return round(score, 2)

def check_vscode_py_indent_docstring__3d465870a880c6e96493a5637b57636e_qw35sft2_9804e750(result, expected, **options):
    """Check indentation (0.5) and docstring presence (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('indent_ok'):
        score += 0.5
    if result.get('has_docstring'):
        score += 0.5
    return score

def check_vscode_multi_settings__8bfa6a8d7119af313a0660a3dea4ceab_qw35sft2_bc343dc4(result, expected, **options):
    """Check multiple VS Code settings with partial credit.

    expected keys:
      checks: list of {"key": str, "value": <expected_value>}
    Scoring: equal weight per check.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    weight = 1.0 / len(checks)
    score = 0.0
    for check in checks:
        key = check['key']
        exp_value = check['value']
        actual = result.get(key)
        if actual == exp_value:
            score += weight
        elif isinstance(exp_value, (int, float)) and isinstance(actual, (int, float)):
            if float(actual) == float(exp_value):
                score += weight
    return min(round(score, 4), 1.0)

def check_vscode_keybinding__70885a60e67c4def20a9d0e3284c1793_qw35sft2_500cb562(result, expected, **options):
    """Check if VS Code keybindings.json contains the expected keybinding entry.

    expected keys:
        key     (str): the keyboard shortcut, e.g. 'ctrl+alt+f'
        command (str): the VS Code command id
    """
    if result.get('error') or not result.get('entries'):
        return 0.0
    entries = result.get('entries', [])
    expected_key = expected.get('key', '').lower().replace(' ', '')
    expected_command = expected.get('command', '')
    for entry in entries:
        if not isinstance(entry, dict):
            continue
        entry_key = entry.get('key', '').lower().replace(' ', '')
        entry_command = entry.get('command', '')
        if entry_command.startswith('-'):
            continue
        if entry_key == expected_key and entry_command == expected_command:
            return 1.0
    return 0.0

def check_vscode_dual_settings__c64cfd038604e0824dbc8809a35f407d_qw35sft2_830dc09c(result, expected, **options):
    """Check two VS Code settings with partial credit (0.5 each).

    result: local path to the downloaded settings.json file (from vm_file getter)
    expected: rules dict already unwrapped by get_rule()
      - setting1_key: str
      - setting1_val: any
      - setting2_key: str
      - setting2_val: any
    """
    import json
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    score = 0.0
    setting1_key = expected.get('setting1_key')
    setting1_val = expected.get('setting1_val')
    setting2_key = expected.get('setting2_key')
    setting2_val = expected.get('setting2_val')
    if setting1_key is not None and settings.get(setting1_key) == setting1_val:
        score += 0.5
    if setting2_key is not None and settings.get(setting2_key) == setting2_val:
        score += 0.5
    return score

def check_vscode_setting__fa93c32e8afd83b84bab5f585818b0fb_qw35sft2_c69a4670(result, expected, **options):
    """Check that a specific VS Code setting matches the expected value."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    setting_key = expected.get('setting_key', '')
    expected_value = expected.get('expected_value')
    if not setting_key:
        return 0.0
    actual_value = result.get(setting_key)
    if isinstance(expected_value, (int, float)) and isinstance(actual_value, (int, float)):
        return 1.0 if actual_value == expected_value else 0.0
    return 1.0 if actual_value == expected_value else 0.0

def check_ext_and_keybinding__44f31fc07357666acac95e9a419ee84f_qw35sft2_dbb9caad(result, expected, **options):
    """Check extension installed (0.5) + specific keybinding present (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_id = expected.get('ext_id', 'undefined_publisher.test')
    if ext_id in result.get('ext_list', ''):
        score += 0.5
    expected_key = expected.get('key', '')
    expected_command = expected.get('command', '')
    keybindings = result.get('keybindings', [])
    if isinstance(keybindings, list):
        for kb in keybindings:
            if not isinstance(kb, dict):
                continue
            kb_key = kb.get('key', '').lower().replace(' ', '')
            kb_cmd = kb.get('command', '')
            if kb_key == expected_key.lower().replace(' ', '') and kb_cmd == expected_command:
                score += 0.5
                break
    return score

def check_vscode_theme_wordwrap__978e7cba9602912bc979b36fdc172116_qw35sft2_b6f8aa43(result, expected, **options):
    """Check colorTheme (0.5) and editor.wordWrap (0.5) with partial credit."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    score = 0.0
    if result.get('colorTheme') == expected.get('colorTheme'):
        score += 0.5
    if result.get('wordWrap') == expected.get('wordWrap'):
        score += 0.5
    return score

def check_workspace_has_folders__f72e9e4f450564b5101c9194862d205b_qw35sft2_548bf84b(result, expected, **options):
    """Check if workspace contains all required folder names. Partial credit per folder."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required = expected.get('required_folder_names', [])
    folder_names = result.get('folder_names', [])
    if not required:
        return 0.0
    per_folder = 1.0 / len(required)
    score = 0.0
    for req in required:
        if req in folder_names:
            score += per_folder
    return min(score, 1.0)

def check_vscode_settings__d45b70bac3694dd736206b10dff26100_qw35sft2_ed066c92(result, expected, **options):
    """Check that VS Code settings.json contains the expected key-value pairs."""
    if not isinstance(result, dict) or 'error' in result:
        return 0.0
    for key, expected_val in expected.items():
        actual_val = result.get(key)
        if actual_val != expected_val:
            return 0.0
    return 1.0

def check_vscode_exclude_pytest__6a66711a30fb269a0e7d77ca8b497959_qw35sft2_fd948bea(result, expected, **options):
    """Check files.exclude has pycache pattern (0.5) and pytest_cache pattern (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    files_exclude = result.get('files_exclude', {})
    pattern1 = expected.get('pattern1', '**/__pycache__')
    pattern2 = expected.get('pattern2', '**/.pytest_cache')
    if files_exclude.get(pattern1):
        score += 0.5
    if files_exclude.get(pattern2):
        score += 0.5
    return score

def check_vscode_remove_and_add_keybinding__8bf5980383e9d80a451eac2b89ccf1df_qw35sft2_30e84c9c(result, expected, **options):
    """Check ctrl+f removal for list.find and ctrl+alt+e addition for Explorer. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('ctrlf_removed'):
        score += 0.5
    if result.get('explorer_shortcut_added'):
        score += 0.5
    return score

def check_vscode_word_wrap_settings__35c11665b6e5c9d9c9abe7f855ccfe46_qw35sft2_91d69975(result, expected, **options):
    """Check wordWrapColumn and wordWrap mode. 0.5 per sub-goal."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_col = result.get('wordWrapColumn')
    exp_col = expected.get('wordWrapColumn')
    if exp_col is not None and actual_col == exp_col:
        score += 0.5
    actual_wrap = result.get('wordWrap')
    exp_wrap = expected.get('wordWrap')
    if exp_wrap is not None and actual_wrap == exp_wrap:
        score += 0.5
    return score

def check_vscode_ext_autosave__e2317900e4eb6fe8b0284d167aa38270_qw35sft2_3c369c77(result, expected, **options):
    """Check extension installed (0.5) and auto-save is enabled in settings (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_list = result.get('extensions', '')
    required_ext = expected.get('required_ext', 'ms-python.python').lower()
    if required_ext in ext_list:
        score += 0.5
    settings = result.get('settings', {})
    auto_save = settings.get('files.autoSave', 'off')
    if auto_save and str(auto_save).lower() != 'off':
        score += 0.5
    return score

def check_workspace_and_extension__bead7b293cda299d38413ae3ff873cc8_qw35sft2_1a5e842f(result, expected, **options):
    """Partial credit: 0.5 for workspace saved, 0.5 for required extension installed."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    required_ext = expected.get('required_extension', 'ms-python.python').lower()
    if result.get('workspace_exists', False):
        score += 0.5
    installed = result.get('installed_extensions', [])
    if any((required_ext in ext for ext in installed)):
        score += 0.5
    return score

def check_vscode_py_function_file__9e0f9c32752df4ad0956b05cabebbc98_qw35sft2_e4880cfc(result, expected, **options):
    """Partial credit: 0.34 for file existing, 0.33 for having a function, 0.33 for correct function name."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_function = expected.get('expected_function_name', '')
    if result.get('exists'):
        score += 0.34
    if result.get('has_function'):
        score += 0.33
        actual_name = result.get('function_name', '')
        if expected_function and actual_name == expected_function:
            score += 0.33
    return min(score, 1.0)

def check_vscode_workspace_in_storage__5752e6d2cd4f13c8d24f076655028c23_qw35sft2_5ec3117b(result, expected, **options):
    """Return 1.0 if the target workspace appears in VS Code's workspaceStorage."""
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('workspace_opened') else 0.0

def check_vscode_py_indent_blankline__2c836d0eca9976b79042c49c51fe71cf_qw35sft2_62042e83(result, expected, **options):
    """Check indentation (0.5) and blank line after function def (0.5) with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('indent_ok'):
        score += 0.5
    if result.get('has_blank_after_def'):
        score += 0.5
    return score

def check_vscode_multi_settings__9ec2182c98e7d1d45bbbb3ae2d0db42e_qw35sft2_feab433d(result, expected, **options):
    """Check multiple VS Code settings with partial credit.

    expected keys:
      checks: list of {"key": str, "value": <expected_value>}
    Scoring: equal weight per check.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    weight = 1.0 / len(checks)
    score = 0.0
    for check in checks:
        key = check['key']
        exp_value = check['value']
        actual = result.get(key)
        if actual == exp_value:
            score += weight
        elif isinstance(exp_value, (int, float)) and isinstance(actual, (int, float)):
            if float(actual) == float(exp_value):
                score += weight
    return min(round(score, 4), 1.0)

def check_vscode_debug_focus_wrap__09cc1b24020e8f7c0f7b48c6a4d96413_qw35sft2_19312b4c(result, expected, **options):
    """Check debug.focusEditorOnBreak is False and editor.wordWrap is 'on'. 0.5 credit each."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('focus_editor_on_break') == expected.get('focus_editor_on_break'):
        score += 0.5
    if result.get('word_wrap') == expected.get('word_wrap'):
        score += 0.5
    return min(score, 1.0)

def check_vscode_setting__a4b8eefcd32d2ac2b259c28c71814550_qw35sft2_8583a75e(result, expected, **options):
    """Check that a specific VS Code setting matches the expected value."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    setting_key = expected.get('setting_key', '')
    expected_value = expected.get('expected_value')
    if not setting_key:
        return 0.0
    actual_value = result.get(setting_key)
    if isinstance(expected_value, (int, float)) and isinstance(actual_value, (int, float)):
        return 1.0 if actual_value == expected_value else 0.0
    return 1.0 if actual_value == expected_value else 0.0

def check_workspace_has_folder__d5491cdef8176b887dce317a57139150_qw35sft2_0784154b(result, expected, **options):
    """Check if workspace contains the required single folder name. Returns 1.0 or 0.0."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    required = expected.get('required_folder_name', '')
    folder_names = result.get('folder_names', [])
    return 1.0 if required and required in folder_names else 0.0

def check_vscode_setting__5cd158cc213aefb25242b372a3862ff3_qw35sft2_e8a19b8e(result, expected, **options):
    """Check if VS Code settings.json has the expected key=value.

    expected keys:
        setting_key   (str): the settings.json key, e.g. 'editor.wordWrap'
        setting_value (any): the expected value, e.g. 'on'
    """
    if result.get('error'):
        return 0.0
    settings = result.get('settings', {})
    expected_key = expected.get('setting_key', '')
    expected_value = expected.get('setting_value')
    if not expected_key:
        return 0.0
    actual_value = settings.get(expected_key)
    return 1.0 if actual_value == expected_value else 0.0

def check_vscode_wrap_and_format__95a613ea8825af09f40057c978f02ad7_qw35sft2_54a5ef43(result, expected, **options):
    """Check wordWrapColumn and formatOnSave. 0.5 per sub-goal."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_col = result.get('wordWrapColumn')
    exp_col = expected.get('wordWrapColumn')
    if exp_col is not None and actual_col == exp_col:
        score += 0.5
    actual_fmt = result.get('formatOnSave')
    exp_fmt = expected.get('formatOnSave')
    if exp_fmt is not None and actual_fmt == exp_fmt:
        score += 0.5
    return score

def check_vscode_locale__215b1cc2c3da9fb70c240f9f4f8df6b5_qw35sft2_c98e144e(result, expected, **options):
    """Check that VS Code argv.json contains the expected locale value.

    result: local file path to downloaded argv.json (string from vm_file getter)
    expected: rules dict (already unwrapped by get_rule()), contains 'expected_locale'
    """
    expected_locale = expected.get('expected_locale', '').lower()
    if not expected_locale:
        return 0.0
    if not result or not os.path.exists(str(result)):
        return 0.0
    try:
        with open(str(result), 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    locale = data.get('locale', '').lower()
    if locale == expected_locale or locale.startswith(expected_locale):
        return 1.0
    return 0.0

def check_workspace_and_wordwrap__1713b5c17b096bf267a520ed66d388f0_qw35sft2_8d1703c8(result, expected, **options):
    """Partial credit: 0.5 for workspace saved, 0.5 for word wrap enabled."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_wrap = expected.get('expected_word_wrap', 'on')
    if result.get('workspace_exists', False):
        score += 0.5
    actual_wrap = result.get('editor_word_wrap')
    if actual_wrap == expected_wrap:
        score += 0.5
    return score

def check_vscode_exclude_autosave__7797789b3c2086d41b9d5f1699195599_qw35sft2_b6097013(result, expected, **options):
    """Check files.exclude has pycache pattern (0.5) and files.autoSave matches (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    files_exclude = result.get('files_exclude', {})
    pycache_pattern = expected.get('pycache_pattern', '**/__pycache__')
    if files_exclude.get(pycache_pattern):
        score += 0.5
    expected_autosave = expected.get('auto_save')
    actual_autosave = result.get('auto_save')
    if expected_autosave is not None and actual_autosave == expected_autosave:
        score += 0.5
    return score

def check_vscode_desktop_file_comment__eea98fddd6e0176c917193edbee2b203_qw35sft2_d3906e8a(result, expected, **options):
    """Partial credit: 0.5 for file existing, 0.5 for first line being a Python comment."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.5
        first_line = result.get('first_line', '')
        if first_line.startswith('#'):
            score += 0.5
    return min(score, 1.0)

def check_vscode_dual_settings__37e5b967fa33d7c5b2d11883c3b2a4f6_qw35sft2_4811c951(result, expected, **options):
    """Check two VS Code settings with partial credit (0.5 each).

    result: local path to the downloaded settings.json file (from vm_file getter)
    expected: rules dict already unwrapped by get_rule()
      - setting1_key: str
      - setting1_val: any
      - setting2_key: str
      - setting2_val: any
    """
    import json
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    score = 0.0
    setting1_key = expected.get('setting1_key')
    setting1_val = expected.get('setting1_val')
    setting2_key = expected.get('setting2_key')
    setting2_val = expected.get('setting2_val')
    if setting1_key is not None and settings.get(setting1_key) == setting1_val:
        score += 0.5
    if setting2_key is not None and settings.get(setting2_key) == setting2_val:
        score += 0.5
    return score

def check_vscode_ext_wordwrap__88ac2a64a7ee030ff707715de339ea34_qw35sft2_5cfa3a00(result, expected, **options):
    """Check extension installed (0.5) and editor.wordWrap is 'on' (0.5)."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    ext_list = result.get('extensions', '')
    required_ext = expected.get('required_ext', 'ms-python.python').lower()
    if required_ext in ext_list:
        score += 0.5
    settings = result.get('settings', {})
    word_wrap = settings.get('editor.wordWrap', 'off')
    if str(word_wrap).lower() == 'on':
        score += 0.5
    return score

def check_vscode_multi_settings__d122d575a22853e47d56ee92b6c18378_qw35sft2_33e4705c(result, expected, **options):
    """Check multiple VS Code settings with partial credit.

    expected keys:
      checks: list of {"key": str, "value": <expected_value>}
    Scoring: equal weight per check.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    weight = 1.0 / len(checks)
    score = 0.0
    for check in checks:
        key = check['key']
        exp_value = check['value']
        actual = result.get(key)
        if actual == exp_value:
            score += weight
        elif isinstance(exp_value, (int, float)) and isinstance(actual, (int, float)):
            if float(actual) == float(exp_value):
                score += weight
    return min(round(score, 4), 1.0)

def check_vscode_open_workspace__c37846152daf06e05c18277c5b1f636d_qw35sft2_caafa045(result, expected, **options):
    """Check if VS Code has the expected workspace folder open.

    expected keys (after get_rule unwrapping):
        expected_workspace_name (str): e.g. "project"
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('workspace_name')
    if actual is None:
        return 0.0
    expected_name = str(expected.get('expected_workspace_name', '')).lower()
    if not expected_name:
        return 0.0
    return 1.0 if actual.lower() == expected_name else 0.0

def check_workspace_folders_replaced__edd7ecbec0169642eef724f574a58652_qw35sft2_3273d2e6(result, expected, **options):
    """Partial-credit check: required folders present + excluded folders absent.
    Equal weight per check across all required + excluded items.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    folder_names = result.get('folder_names', [])
    required = expected.get('required_folder_names', [])
    excluded = expected.get('excluded_folder_names', [])
    total_checks = len(required) + len(excluded)
    if total_checks == 0:
        return 0.0
    per_check = 1.0 / total_checks
    score = 0.0
    for req in required:
        if req in folder_names:
            score += per_check
    for exc in excluded:
        if exc not in folder_names:
            score += per_check
    return min(round(score, 4), 1.0)

def vscode_workspace_saved__f91feeeb0a889e0b6b274960d3d06c7d_qw35sft2_280f3ee8(result, expected, **options):
    """Return 1.0 if workspace file exists and is valid JSON at /home/user/project.code-workspace."""
    if not isinstance(result, dict):
        return 0.0
    if not result.get('workspace_exists', False):
        return 0.0
    if not result.get('valid_json', False):
        return 0.5
    return 1.0

def check_vscode_tab_size__015bdcf4a8f96215b5579fc53db2d526_qw35sft2_451e5295(result, expected, **options):
    """Check editor.tabSize equals expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('tabSize')
    exp = expected.get('tabSize')
    if exp is not None and actual == exp:
        return 1.0
    return 0.0

def check_vscode_debug_focus_terminal__29baecbd3036f18eff818341b1c0e4e8_qw35sft2_814b6d3c(result, expected, **options):
    """Check debug.focusEditorOnBreak is False and terminal.integrated.cursorBlinking is True. 0.5 credit each."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('focus_editor_on_break') == expected.get('focus_editor_on_break'):
        score += 0.5
    if result.get('cursor_blinking') == expected.get('cursor_blinking'):
        score += 0.5
    return min(score, 1.0)

def check_vscode_setting__d97bffb72b4b8c27def4b087597b01a0_qw35sft2_f7177a16(result, expected, **options):
    """Check that a specific VS Code setting matches the expected value."""
    if isinstance(result, dict) and 'error' in result:
        return 0.0
    setting_key = expected.get('setting_key', '')
    expected_value = expected.get('expected_value')
    if not setting_key:
        return 0.0
    actual_value = result.get(setting_key)
    if isinstance(expected_value, (int, float)) and isinstance(actual_value, (int, float)):
        return 1.0 if actual_value == expected_value else 0.0
    return 1.0 if actual_value == expected_value else 0.0

def check_vscode_dual_settings__b75b59f2415937a4444d85150c1ab46c_qw35sft2_5a8c7a44(result, expected, **options):
    """Check two VS Code settings with partial credit (0.5 each).

    result: local path to the downloaded settings.json file (from vm_file getter)
    expected: rules dict already unwrapped by get_rule()
      - setting1_key: str
      - setting1_val: any
      - setting2_key: str
      - setting2_val: any
    """
    import json
    if not result or not isinstance(result, str):
        return 0.0
    try:
        with open(result, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    score = 0.0
    setting1_key = expected.get('setting1_key')
    setting1_val = expected.get('setting1_val')
    setting2_key = expected.get('setting2_key')
    setting2_val = expected.get('setting2_val')
    if setting1_key is not None and settings.get(setting1_key) == setting1_val:
        score += 0.5
    if setting2_key is not None and settings.get(setting2_key) == setting2_val:
        score += 0.5
    return score

def check_vscode_single_ext__521eb992bf97de2d79e4aa6c081db684_qw35sft2_ba37d45a(result, expected, **options):
    """Check if a single VS Code extension is installed."""
    if not isinstance(result, dict):
        return 0.0
    ext_list = result.get('extensions', '')
    required_ext = expected.get('required_ext', '').lower()
    if required_ext and required_ext in ext_list:
        return 1.0
    return 0.0

def check_vscode_active_file__6a582969bc6d0b9624887969279510e3_qw35sft2_33b62156(result, expected, **options):
    """Check if the expected file is the active editor in VS Code.

    expected keys (after get_rule unwrapping):
        expected_file (str): filename to check, e.g. "main.py"
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('active_file')
    if actual is None:
        return 0.0
    expected_file = str(expected.get('expected_file', ''))
    if not expected_file:
        return 0.0
    return 1.0 if actual.lower() == expected_file.lower() else 0.0

def check_vscode_multi_settings__7e25498cca269e458737cc32ebf8beb0_qw35sft2_9903d5f1(result, expected, **options):
    """Check multiple VS Code settings with partial credit.

    expected keys:
      checks: list of {"key": str, "value": <expected_value>}
    Scoring: equal weight per check.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    weight = 1.0 / len(checks)
    score = 0.0
    for check in checks:
        key = check['key']
        exp_value = check['value']
        actual = result.get(key)
        if actual == exp_value:
            score += weight
        elif isinstance(exp_value, (int, float)) and isinstance(actual, (int, float)):
            if float(actual) == float(exp_value):
                score += weight
    return min(round(score, 4), 1.0)
