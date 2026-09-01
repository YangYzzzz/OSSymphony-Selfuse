"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_impress_save_settings__a73132a8bbdfe6ee96f72881e60f8547', 'check_slide_hidden__266b91e28f8289ba603f3990e8aa8eb6', 'check_pptx_bold__c8d3ea84bab1df08b4dcc2adb9177ac5', 'check_pptx_textbox_font_sizes__a0ffb78fbe8b896461c46fcb419a2efb', 'check_new_slide_added__964158cb23267c07918d09a84fcd28cc', 'check_slide_title_text__6c6be158040f1b190133813bff96d3a4', 'check_pptx_slide_title__36872752e18a684ac0b72ebb0b9e2c33', 'check_pptx_slide_subtitle__b0063e78116208a0699a2e91287a92ca', 'check_pptx_title_format__322bfa5ed10c48cee5b1ded75b68bcd6', 'check_impress_title_color__9e6f539665d909e9d4b70212d289e0ee', 'check_pptx_orientation__2f288a1714ed5a4b53fed9b568a61679', 'check_pptx_title_font_color__3536f82770e6c6d12bb7de4bc43120d5', 'check_slide_audio__1b19ca869f4a35a02b8d110c7e17b7e9', 'check_pptx_textbox_font_sizes__8e211b70307c1527c1dc374c2a15baa9', 'check_pptx_slide_title_and_bg__93a38408d9330c4f554c19c28759d1cc', 'check_pptx_shape_dimensions__30ab1f17976bddf578367b54740da469', 'check_impress_strikethrough__0be9a8c316c2b4d9a9ba4471962b751f', 'check_slide1_italic_red__7d9f2200326202f628d33a5691960f9c', 'check_impress_strikethrough__bbd5669ff0c1c503844620466c04ffda', 'check_slide_notes__4bc42c00ac052a9214daa0b289a72718', 'check_impress_slide2_state__33bd53bce176d3722b0b7a880cb5489e', 'check_slide_has_bg_image__54b7105f00ae5aefd1a4f4e9a1686c3b', 'check_impress_text_and_bg__6d65cdf6ba3aad0bcf7dd7e1b9c52c8f', 'check_pptx_shape_dimensions__dc32e28ff1403085a7367144a8c3b449', 'check_pptx_slide_bg_color__6efa38f597117d6ec1022844076cbdb0', 'check_pptx_slide_title__95b2513b7f5b16d49b143853ab08b964', 'check_slide_orientation__662bf4c59a75818f259c7e0aef709493', 'check_slide_title_font__05b860587b204cf23347a8d6b2a21320', 'check_slide_deleted__b3862aa49d8968cf4ed21032ba8b229f', 'check_impress_new_slide__3a9a25804cb58719a25a23e52e7e124a', 'check_impress_table_with_content__608b93868302b656e50a31095c9d18e9', 'check_slide_orientation__326a96d7a18db487b955505da7fa674b', 'check_slide_title__ebdab037a8eab6c57e56f5796ba7c0ab', 'check_pptx_pic_height_and_fonts__5b53387f975ab04f84e7061a6b89939a', 'check_impress_slide_deletion__1abdd5c6fe686f66f2606ee0a55bf0d1', 'check_impress_title_and_bg__ac120f719f06d988c1c09c7022ee4b50', 'check_pptx_image_size__19c7a30018144cafc9f4ebd4b22910ef', 'check_pptx_slide_count__f16a19b8aecc56a3153baf7aa9668d95', 'check_pptx_title_text__661ea64646d6a38833b5fdf5450470f8', 'check_impress_all_title_fonts__9c7bd04930deae5cddd57fa6475996f4', 'check_pptx_pic_height_and_fonts__4eb9ebb1f165db3d3eb60748b0824590', 'check_snapshot_and_slide_bg__da1a5547fec970e705ee94d1f2bd1e91', 'check_impress_strikethrough__941b469f6d56a3f2848b88f32a6979cb', 'check_pptx_orientation__f921e5b58a724c7baa4b59415c1c4019', 'check_slide_title_font_color__66374a642830d11922d74321b71d53a4', 'check_slide_dimensions__54270a3c8d106e03ce2d0c72089a3417', 'check_pptx_slide_title__24e166a4a3628d77570c2919341fa3d0', 'check_presentation_summary__f7a75c940bb0ba7687c094e66be96f16', 'check_impress_title_fontsize__f67d7e69d6a85882636e0a41dabd5c89', 'check_pptx_slide_count_and_title__e6d2505b7a2f1323c83204293d9b2ff3', 'check_slide_title_text__f01c4fcc320219e33dd1b3dfc752a5b6', 'check_slide_count__55d37620fcf7998461b41e2b1579f16b', 'check_slide_notes__44fd5b3a7c457950926c4afa82f02524', 'check_slide2_underline_48pt__203dc282bb468063e328a7afef0f556d', 'check_pptx_pic_width_and_fonts__1a124508c5c52deecd8c9cfabb1cdd16', 'check_slide_text__64fdb0e52d9e147a698cbe6791191c81', 'check_slide_title__7b96cb41298a29a8e5448750424ab134', 'check_slide_title__94c866527f562343eb1ed6d961346913', 'check_slide5_bold_36pt__e80588656c5c992160d4d54a69813615', 'check_slide_bg_color__147ec12b6038f2d22f20700395a41639', 'check_slide_title__497f89854667fa402a3b0a0e21fc061d', 'check_slide_title__07d639643e02063bb060c675d5eb936b', 'check_impress_slide_duplicate__b84f4f9ed96e0e954dfada4f0a280734', 'check_slide_titles__d445ea0218043b2b599f555fb0b0c923', 'check_pptx_title_bold__7728836b00929340fdef9a2a6602ef09', 'check_impress_slide_table__dabff5cb0982227a6c01462d5c667c85', 'check_pptx_paragraph_contains__10c4d6d7b532dd358d2936a84f08db4a', 'check_pptx_image_size__32c4d5ced54cde2e00ee07b6ba3973e4', 'check_pptx_all_slides_bg_color__88a6aaaf9999692a07033829e3809d82', 'check_pptx_textbox_font_sizes__d842f0d6650bc2853fbd82b34d9bb662', 'check_impress_title_font_props__8c5c9bdd88c0ba0d1f43df42a252cfae', 'check_pptx_shape_dimensions__37468b49cfdc53eaca3dfbbee86985a3', 'check_impress_bg_and_subtitle__024a68353adff6dd2d6a01cc842a46ed', 'check_pptx_transitions__a8621b66c753635f21d8c612a865d275', 'check_slide_image_inserted__e8973b7aba53f626fa037406b1da4c8e', 'check_impress_title_bold__6e02436ff918c1445d486be48242a3ca', 'check_pptx_multi_slide_images__e6e9ce905cbd4fd74addd1fc2cd640de', 'check_slide_bg_color__1dbf2adcf932ab2d0e517ac3d78f9d08', 'check_audio_on_slide__5688a332db768af3543ab15e11555c76', 'check_slide_count__3e661325ce185501287a57bdd6c09b52', 'check_slide_count_blank__c9d9670d65b42979a3ba7969e086139a', 'check_pptx_italic__979792a60bf386d8e1c5a4d702d19829', 'check_impress_slide_reorder__88930a069947104f21bb2b61fdebcc2c', 'check_pptx_title_text__927f9acbea3f73fbaa9e281542feb428', 'check_slide_title_and_color__fd2e510c274f3f89d65383d731b18828', 'check_slide_count__1573725569136f71d4ad79ac5ffbf35a', 'check_slide_duplicated__e9203909bdddcf8283496218733a0b9b', 'check_pptx_slide_count__d96a0661e83ae9239b523060c49e36f2_qw35sft2_4679f70e', 'check_pptx_table_row0__27882b2f85f9d8345d5e8153a8f07379_qw35sft2_bf84994b', 'check_pptx_summary_notes__f1eef050d638a5408c254c4e620a5302_qw35sft2_a99f5872', 'check_pptx_last_moved_to_first__b9dba6723fe45eeec63ed20d4f46d22d_qw35sft2_376a34c4', 'check_impress_title_bottom_and_text__dea2c3a0c75be4ec786c2a7bbd9ea59f_qw35sft2_ff64be5a', 'check_pptx_slide_transitions__f348aaf139c455284fcbda7cfe577da5_qw35sft2_a8a2fc6d', 'check_pptx_image_size__7a22b4a6c189fa88dd4ad54ac62a04cd_qw35sft2_5933753b', 'check_pptx_bg_solid__7f279c8dbe68c5780c5b5f60310b951b_qw35sft2_a43b9793', 'check_all_slides_fade__85331dc0cf4d821fa2bad3e27657d683_qw35sft2_3c31d6ea', 'check_pptx_slide_bg_color__1a5b2a2b44ee739e3486ddbc20ac1c87_qw35sft2_36618094', 'check_pptx_slide_has_solid_bg__30ec500dc9631258b0ce307c50b98145_qw35sft2_6aa292db', 'check_impress_slide4_contact_shapes__5eae632a29df23abd3f96eee0032e6dc_qw35sft2_3d109fc7', 'check_impress_underline__fa317ec6034191f7b02a906f78ff4b09_qw35sft2_7c170663', 'check_impress_audio_slide1_trans_slide3__1b01b098ee4b6dfab8bc1b2d24723a1f_qw35sft2_f2414360', 'check_blank_slide_count__da57195bec8b00f5d49d80dfe28b0049_qw35sft2_946e4f58', 'check_impress_bullet_newpara__12c621423b2418d27e86fd746731530a_qw35sft2_85f3aeeb', 'check_pptx_text_alignments__6253f908caa7dec8dc9ccbad6c2f7332_qw35sft2_756004b9', 'check_pptx_slide_font_size__4b92c85bf54981a99e31cdf8daa555a4_qw35sft2_3d7ddbc4', 'check_pptx_slide3_text_colors__9e2b961f79a8b89da2fe1a7f40c57f65_qw35sft2_0a01b8df', 'check_pptx_slide3_para_texts__c0cb167c6059fd8f9f81e328b9599144_qw35sft2_a0cf6b20', 'check_impress_multi_title_color__14718ed4f8dc81749072a781b972112d_qw35sft2_e6acc344', 'check_impress_slide1_title__992b8b598a921e5736a852e261be2237_qw35sft2_e1bfdb4f', 'check_pptx_slide2_title_state__1225c3617161e6b819a6b6275a5c9749_qw35sft2_d3d81117', 'impress_slide3_table_at_bottom__db4d3e235d7b0e304c073a25921cbfa2_qw35sft2_ca53a561', 'check_slide14_font_sizes__ec3f2f9f447d657b0668843895ae7804_qw35sft2_b758be1f', 'check_impress_notes_bg_title__cb759dd1cf89c2aec2802f8543d62400_qw35sft2_1d03f915', 'check_pptx_text_color__39fe73dd35955e4e54b2d721092898e1_qw35sft2_d603cba2', 'check_pptx_text_color__c830024b32b30956e96002af134f6f4a_qw35sft2_a9a956ae', 'check_slide2_bg_subtitle_color__2d907c4b088ee7e1bdf99fbd932517ff_qw35sft2_3aba4d69', 'check_pptx_title_subtitle_font__761afd0460bdb180560f4c7116933937_qw35sft2_8b84cf99', 'check_impress_pptx_props__ed92c6da7770c48f6ebca23c6070cbf5_qw35sft2_55ab6274', 'check_pptx_content_and_table__7f43e77e7e3ce83ed182df2ac03d3d3a_qw35sft2_df76f63e', 'check_picture_size_slide6__f55d52551229ea682c89b1a15474ccff_qw35sft2_54cc8e8d', 'check_impress_panel_slide_count__ddf9714ddb153a97bbdfb4350e733817_qw35sft2_992c6a6b', 'check_pptx_font_size__f00d8a2d9828be4b94677a9606cd7f47_qw35sft2_09ce01c6', 'check_pptx_portrait__2491631e79810e0b815ac32866fd3274_qw35sft2_37505aa0', 'check_pptx_orientation_state__f5c58fc75d059d716d7a9ec9e8e3967d_qw35sft2_4344dc80', 'check_impress_title_at_center__ff47db79abc78d9904bd7e17dbb1e3d3_qw35sft2_e46633fe', 'check_pptx_notes_nonempty__d897409e855e2cbbb791ebde03eeef8a_qw35sft2_6c59bb80', 'check_pptx_slide_title_text__c928e38a883d919c295367c551657155_qw35sft2_48305e8b', 'check_pptx_image_and_transition__4195ceef20cb2389eadfdf0daf18e8e7_qw35sft2_300afbf0', 'check_slide_bg_color__b309be1bc27dcbe28bf5749a28cc9b6f_qw35sft2_04bc853b', 'check_pptx_slide_subtitle_text__4d840f044cdf667c54faf806c6af73e1_qw35sft2_c3b1a416', 'check_pptx_last_two_duplicated__ae132e5d53d956fce4051acdf425fcba_qw35sft2_371de5a2', 'check_blank_slide_count__881be08815c16b0f7ba88245e2d69e10_qw35sft2_c7683f4c', 'check_impress_audio_and_transition__1b19ca869f4a35a02b8d110c7e17b7e9_qw35sft2_f096de71', 'check_pptx_slide_font_italic__5e9f703dc627e70a3b4d09452ec3d3ab_qw35sft2_23ac8c49', 'check_pptx_slides35_colors__0fbb4b3515e1009c69bc08919cd2b491_qw35sft2_82a32b53', 'check_pptx_slide3_subpoint_level__3c50fb39e50d5c4bd155f8f4b7e911d1_qw35sft2_cdee0fb6', 'check_save_and_slide_count__e4d3506bfaa42b2028f1a93213c30bc6_qw35sft2_73a87272', 'check_impress_bullet_underline__396e098239a8be5c53a96d982334cb7b_qw35sft2_6b235811', 'check_impress_options_combo__43da862b24de1bb5002a875bb93d754e_qw35sft2_6c066c1d', 'check_pptx_slide_orientation__1830f729a2ae0637260a8bca58b0c8cd_qw35sft2_c6ecb7c1', 'check_impress_img_top__edae955a62ba6a9d5bff7ce774e3ee10_qw35sft2_5c366587', 'check_impress_title_color_underline__84ce98fac5a701a1ff43f1337a4f838c_qw35sft2_b63365eb', 'check_pptx_slide2_title_bold__ba024eb57784584ce2b08e31db5871ac_qw35sft2_6b563f6f', 'check_impress_slide3_title_bold__7dcab435f4cba9d9e23699eed9a4d408_qw35sft2_1080995a', 'check_pptx_text_alignments__14ff9da6ca84bfee11448e7c7f8efafc_qw35sft2_86d84f40', 'check_impress_notes_bg_transition__5cb083f40f95f75316141f060ee8a1cc_qw35sft2_5f8249e7', 'check_slide14_font_bold__a28710b7c875e09d23efa44313d2c747_qw35sft2_644b3032', 'impress_slide3_bottom_and_title__80d9e8d7efd024cbc51991913269b5fb_qw35sft2_fba30bc1', 'check_impress_slide_title_full__9620e3371e8842ebb86c9c7ef6cb4f00_qw35sft2_fa068e33', 'check_pptx_title_font_transition__7950498f55dec024fcd49c95968dd860_qw35sft2_2c7c0d01', 'check_pptx_text_color__f2292754a9e180ffad3521e45a75c905_qw35sft2_dc1d29ca', 'check_slide2_bg_title_bold__3ee304ba526b814e3d4d6c2be946eae0_qw35sft2_0f122ed4', 'check_pptx_text_color__c8d6e08f7da191ce71fc333f5a053d63_qw35sft2_1bd747b0', 'check_pptx_content_bold__87b69c70091c302ea2d6ddb5f5d9c002_qw35sft2_5ef7630d', 'check_impress_slide_pane__0dbd76d4f35e1d6b60aa20d8163a50c2_qw35sft2_d7e6c0f9', 'check_pptx_font_size__8cd3b2c6f2a2d80e7aa650d3b0962695_qw35sft2_700a1334', 'check_impress_pptx_props__e526a93f073489d810264dc88333d90d_qw35sft2_f459912c', 'check_pptx_has_transition__24d25146d014952eed8e287afb1ec612_qw35sft2_3e97889e', 'check_pptx_transitions__b146bd229b2aa1513d722fe51dc492dc_qw35sft2_1a2a7e26', 'check_pptx_table_row0__a30c5776ecd0bc922b282656ddbb8d0e_qw35sft2_978d7b1e', 'check_pptx_slide_title_text__a33045396be2560e0d47268bd7a68b1f_qw35sft2_af87e143', 'check_slide_bg_color__3e8a909d4fdee8dfddabe4251505a445_qw35sft2_3ed277f0', 'check_blank_slide_count__b298657a7c2d705873344441c4e6373c_qw35sft2_6a7840bc', 'check_pptx_image_and_title__68b16c8fa261fe6f7f7bf99cd53e207c_qw35sft2_a9adceaf', 'check_impress_title_bottom_and_font__8ff66c4d1cd7e4775354b5fdf25e12eb_qw35sft2_d3fefb96', 'check_pptx_notes_title__e6eecf5158b3e658cb4aa48b9b009898_qw35sft2_e6001705', 'check_two_slide_tables__422453ea9133a41f5d3b73c5d61c25f4_qw35sft2_1314e343', 'check_pptx_slide5_text_colors__7491f6c65f733a38f893bdf8cad068b0_qw35sft2_659f48ec', 'check_impress_has_audio__5688a332db768af3543ab15e11555c76_qw35sft2_138ab214', 'check_pptx_multi_slide_bg_colors__01d7b3b210af327bea300a303c9bddd6_qw35sft2_9c2a1925', 'check_impress_video_docs_vlc__bae283564a19c9e043588bc4818877d3_qw35sft2_c7d97e82', 'check_impress_compose4__c60351f10d67b0cf97ff8e001ca27d87_qw35sft2_9a3a7691', 'check_pptx_slide2_bg_and_count__24ae1553823b67948a98514a9ed74ad3_qw35sft2_53ba711c', 'check_pptx_notes_and_docx__1f2059f6b413a703d90fb8946ec0e698_qw35sft2_3c27e232', 'check_pptx_slide2_bg__d5c6210d63b8dbe359320c4bfae3310a_qw35sft2_7e1c3639', 'check_impress_compose0__9880fcd0ac1db0e1d57c564194f1780a_qw35sft2_0c53865a', 'check_vlc_dark_slider_oneinstance__d71f004306b232a878090dc93cee8c3a_qw35sft2_07da77f7', 'check_vlc_dark_slider_minimal__d20b1f499db3b1fe73151927397d72b8_qw35sft2_184f1657', 'check_vlc_dark_slider_maxvol__125763cfb25260df78f204fa214a8b29_qw35sft2_dd32e44e', 'check_vlc_dark_slider_bgcone__914b13738d8dcf766e175adf76980e8b_qw35sft2_0920035b', 'check_vlc_dark_slider_globalkey__c70b3a78c573774afb864d5784b3c8df_qw35sft2_56f9e06c']

def check_impress_save_settings__a73132a8bbdfe6ee96f72881e60f8547(result, expected, **options):
    """Check auto-save interval and backup copy settings with partial credit.

    Scoring:
    - 0.5 for correct auto-save interval (enabled + correct minutes)
    - 0.5 for backup copy enabled
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_minutes = expected.get('expected_minutes')
    if result.get('auto_save_enabled') and result.get('auto_save_minutes') == expected_minutes:
        score += 0.5
    expected_backup = expected.get('expected_backup', True)
    if result.get('backup_enabled') == expected_backup:
        score += 0.5
    return min(score, 1.0)

def check_slide_hidden__266b91e28f8289ba603f3990e8aa8eb6(result, expected, **options):
    """Check that a specific slide is hidden."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_hidden = expected.get('expected_hidden', True)
    expected_title = expected.get('expected_title', 'QUIZ - FORESTS')
    is_hidden = result.get('is_hidden', False)
    slide_title = result.get('slide_title', '')
    slide_count = result.get('slide_count', 0)
    if slide_count == expected.get('expected_slide_count', 12):
        score += 0.4
    if expected_title and slide_title:
        if expected_title.upper() in slide_title.upper():
            score += 0.1
    if is_hidden == expected_hidden:
        score += 0.5
    return min(score, 1.0)

def check_pptx_bold__c8d3ea84bab1df08b4dcc2adb9177ac5(result, expected, **options):
    """Check if text on specified slides has bold formatting. Partial credit per slide."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    slide_indices = expected.get('slide_indices', [0, 1])
    score = 0.0
    weight = 1.0 / len(slide_indices) if slide_indices else 0.0
    for idx in slide_indices:
        key = f'slide_{idx}'
        slide_data = result.get(key, {})
        if slide_data.get('error'):
            continue
        total = slide_data.get('total_runs', 0)
        bold = slide_data.get('bold_runs', 0)
        if total > 0 and bold == total:
            score += weight
    return min(score, 1.0)

def check_pptx_textbox_font_sizes__a0ffb78fbe8b896461c46fcb419a2efb(result, expected, **options):
    """Check if textbox font sizes match expected values. Supports partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    font_sizes = result.get('font_sizes', {})
    expected_sizes = expected.get('expected_sizes', {})
    if not expected_sizes:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_sizes)
    for (text, expected_size) in expected_sizes.items():
        actual_size = font_sizes.get(text)
        if actual_size is not None and abs(actual_size - expected_size) < 0.5:
            score += per_item
    return min(score, 1.0)

def check_new_slide_added__964158cb23267c07918d09a84fcd28cc(result, expected, **options):
    """Check if a new slide was added with the correct title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 29)
    expected_title = expected.get('expected_last_title', '')
    actual_count = result.get('slide_count', 0)
    actual_title = result.get('last_slide_title', '')
    if actual_count >= expected_count:
        score += 0.5
    if actual_title and expected_title:
        if actual_title.strip().lower() == expected_title.strip().lower():
            score += 0.5
        elif expected_title.strip().lower() in actual_title.strip().lower():
            score += 0.25
    return min(score, 1.0)

def check_slide_title_text__6c6be158040f1b190133813bff96d3a4(result, expected, **options):
    """Check if slide title text matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if not expected_title:
        return 0.0
    if actual_title.lower() == expected_title.lower():
        return 1.0
    return 0.0

def check_pptx_slide_title__36872752e18a684ac0b72ebb0b9e2c33(result, expected, **options):
    """Check if the slide title matches expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title_text', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if actual_title == expected_title:
        return 1.0
    if actual_title.lower() == expected_title.lower():
        return 0.8
    return 0.0

def check_pptx_slide_subtitle__b0063e78116208a0699a2e91287a92ca(result, expected, **options):
    """Check if slide subtitle matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('subtitle', '')
    expected_text = expected.get('expected_subtitle', '')
    if actual is None:
        return 0.0
    if actual.strip() == expected_text.strip():
        return 1.0
    return 0.0

def check_pptx_title_format__322bfa5ed10c48cee5b1ded75b68bcd6(result, expected, **options):
    """Check if the title font formatting matches expected values (bold + color)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('bold', False)
    if actual_bold == expected_bold:
        score += 0.5
    expected_color = expected.get('expected_color_rgb', '').upper()
    actual_color = (result.get('font_color_rgb') or '').upper()
    if expected_color and actual_color == expected_color:
        score += 0.5
    return min(score, 1.0)

def check_impress_title_color__9e6f539665d909e9d4b70212d289e0ee(result, expected, **options):
    """Check if title font color matches expected color.

    Uses fuzzy matching with tolerance for color distance.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_hex = result.get('color_rgb')
    expected_hex = expected.get('expected_color_rgb')
    if actual_hex is None or expected_hex is None:
        return 0.0
    try:
        actual_hex = actual_hex.lstrip('#')
        expected_hex = expected_hex.lstrip('#')
        ar = int(actual_hex[0:2], 16)
        ag = int(actual_hex[2:4], 16)
        ab = int(actual_hex[4:6], 16)
        er = int(expected_hex[0:2], 16)
        eg = int(expected_hex[2:4], 16)
        eb = int(expected_hex[4:6], 16)
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        threshold = expected.get('color_tolerance', 30)
        if distance <= threshold:
            return 1.0
        return 0.0
    except (ValueError, IndexError):
        return 0.0

def check_pptx_orientation__2f288a1714ed5a4b53fed9b568a61679(result, expected, **options):
    """Check if presentation has expected orientation."""
    if isinstance(result, str) or not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_orientation = expected.get('orientation', 'portrait').lower()
    actual_orientation = result.get('orientation', '').lower()
    return 1.0 if actual_orientation == expected_orientation else 0.0

def check_pptx_title_font_color__3536f82770e6c6d12bb7de4bc43120d5(result, expected, **options):
    """Check if the title font color matches the expected color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_color = expected.get('expected_color', '000000')
    first_color = result.get('first_color')
    if first_color is None:
        return 0.0
    if first_color.upper() == expected_color.upper():
        return 1.0
    return 0.0

def check_slide_audio__1b19ca869f4a35a02b8d110c7e17b7e9(result, expected, **options):
    """Check if the correct audio was inserted on the specified slide.
    Partial credit: 0.5 for having audio on the slide, 0.5 for matching the reference.
    """
    if isinstance(result, dict) and result.get('error') and (not result.get('has_audio')):
        return 0.0
    score = 0.0
    if result.get('has_audio'):
        score += 0.5
    if result.get('audio_match'):
        score += 0.5
    return score

def check_pptx_textbox_font_sizes__8e211b70307c1527c1dc374c2a15baa9(result, expected, **options):
    """Check if textbox font sizes match expected values. Supports partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    font_sizes = result.get('font_sizes', {})
    expected_sizes = expected.get('expected_sizes', {})
    if not expected_sizes:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_sizes)
    for (text, expected_size) in expected_sizes.items():
        actual_size = font_sizes.get(text)
        if actual_size is not None and abs(actual_size - expected_size) < 0.5:
            score += per_item
    return min(score, 1.0)

def check_pptx_slide_title_and_bg__93a38408d9330c4f554c19c28759d1cc(result, expected, **options):
    """Check title text and background color with partial credit (0.5 each)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title_text', '')
    if actual_title and expected_title:
        if actual_title.strip().lower() == expected_title.strip().lower():
            score += 0.5
    expected_color = expected.get('expected_bg_color', '').upper()
    actual_color = result.get('bg_color')
    if actual_color and expected_color:
        actual_color = actual_color.upper()
        if actual_color == expected_color:
            score += 0.5
        else:
            tolerance = expected.get('tolerance', 30)
            try:
                ar = int(actual_color[0:2], 16)
                ag = int(actual_color[2:4], 16)
                ab = int(actual_color[4:6], 16)
                er = int(expected_color[0:2], 16)
                eg = int(expected_color[2:4], 16)
                eb = int(expected_color[4:6], 16)
                distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
                if distance <= tolerance:
                    score += 0.5
            except (ValueError, IndexError):
                pass
    return score

def check_pptx_shape_dimensions__30ab1f17976bddf578367b54740da469(result, expected, **options):
    """Check if shape dimensions match expected values with tolerance and partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    tolerance = expected.get('tolerance', 0.5)
    score = 0.0
    weight = 1.0 / len(checks)
    for check in checks:
        key = check['key']
        expected_value = check['value']
        actual = result.get(key)
        if actual is not None and abs(actual - expected_value) <= tolerance:
            score += weight
    return min(round(score, 4), 1.0)

def check_impress_strikethrough__0be9a8c316c2b4d9a9ba4471962b751f(result, expected, **options):
    """Check if target paragraphs have strikethrough formatting."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    target_texts = expected.get('target_texts', [])
    if not target_texts:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(target_texts)
    for target in target_texts:
        for para in paragraphs:
            if target.lower() in para['text'].lower():
                if para['strikethrough']:
                    score += per_item
                break
    return min(score, 1.0)

def check_slide1_italic_red__7d9f2200326202f628d33a5691960f9c(result, expected, **options):
    """Check that title on slide 1 is italic and red."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    runs = result.get('runs', [])
    if not runs:
        return 0.0
    score = 0.0
    all_italic = all((r.get('italic') is True for r in runs))
    if all_italic:
        score += 0.5
    expected_color = expected.get('expected_color', 'FF0000')
    all_red = all((r.get('color_rgb', '').upper() == expected_color.upper() for r in runs))
    if all_red:
        score += 0.5
    return score

def check_impress_strikethrough__bbd5669ff0c1c503844620466c04ffda(result, expected, **options):
    """Check if target paragraphs have strikethrough formatting."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    target_texts = expected.get('target_texts', [])
    if not target_texts:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(target_texts)
    for target in target_texts:
        for para in paragraphs:
            if target.lower() in para['text'].lower():
                if para['strikethrough']:
                    score += per_item
                break
    return min(score, 1.0)

def check_slide_notes__4bc42c00ac052a9214daa0b289a72718(result, expected, **options):
    """Check if slide notes contain the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    notes_text = result.get('notes_text', '')
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    if expected_text.lower() in notes_text.lower():
        return 1.0
    return 0.0

def check_impress_slide2_state__33bd53bce176d3722b0b7a880cb5489e(result, expected, **options):
    """Check slide 2 title underline and body text content.

    Partial credit: 0.5 for underline, 0.5 for body text.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_underline = expected.get('expected_underline', True)
    actual_underline = result.get('title_underline')
    if actual_underline == expected_underline:
        score += 0.5
    expected_body = expected.get('expected_body_text', '')
    actual_body = result.get('body_text', '') or ''
    if expected_body and expected_body.strip().lower() in actual_body.strip().lower():
        score += 0.5
    return min(score, 1.0)

def check_slide_has_bg_image__54b7105f00ae5aefd1a4f4e9a1686c3b(result, expected, **options):
    """Check if the slide has a background image set."""
    if isinstance(result, dict) and result.get('has_bg_image', False):
        return 1.0
    return 0.0

def check_impress_text_and_bg__6d65cdf6ba3aad0bcf7dd7e1b9c52c8f(result, expected, **options):
    """Check slide 6 text and slide 3 background color.
    Partial credit: 0.5 for correct text, 0.5 for correct background.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', 'Thank You!')
    actual_text = result.get('slide6_text', '')
    if actual_text and expected_text.lower().strip('!').strip() in actual_text.lower():
        score += 0.5
    expected_color = expected.get('expected_bg_color', '90EE90').upper()
    actual_color = (result.get('slide3_bg_color') or '').upper()
    if actual_color:
        try:
            er = int(expected_color[0:2], 16)
            eg = int(expected_color[2:4], 16)
            eb = int(expected_color[4:6], 16)
            ar = int(actual_color[0:2], 16)
            ag = int(actual_color[2:4], 16)
            ab = int(actual_color[4:6], 16)
            tolerance = expected.get('color_tolerance', 60)
            if abs(er - ar) <= tolerance and abs(eg - ag) <= tolerance and (abs(eb - ab) <= tolerance):
                score += 0.5
        except (ValueError, IndexError):
            pass
    return min(score, 1.0)

def check_pptx_shape_dimensions__dc32e28ff1403085a7367144a8c3b449(result, expected, **options):
    """Check if shape dimensions match expected values with tolerance and partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    tolerance = expected.get('tolerance', 0.5)
    score = 0.0
    weight = 1.0 / len(checks)
    for check in checks:
        key = check['key']
        expected_value = check['value']
        actual = result.get(key)
        if actual is not None and abs(actual - expected_value) <= tolerance:
            score += weight
    return min(round(score, 4), 1.0)

def check_pptx_slide_bg_color__6efa38f597117d6ec1022844076cbdb0(result, expected, **options):
    """Check if slide background matches expected color within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    color_rgb = result.get('color_rgb')
    if color_rgb is None:
        return 0.0
    expected_color = expected.get('expected_color', '').upper()
    actual_color = str(color_rgb).upper()
    if actual_color == expected_color:
        return 1.0
    tolerance = expected.get('tolerance', 30)
    try:
        ar = int(actual_color[0:2], 16)
        ag = int(actual_color[2:4], 16)
        ab = int(actual_color[4:6], 16)
        er = int(expected_color[0:2], 16)
        eg = int(expected_color[2:4], 16)
        eb = int(expected_color[4:6], 16)
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        if distance <= tolerance:
            return 1.0
    except (ValueError, IndexError):
        pass
    return 0.0

def check_pptx_slide_title__95b2513b7f5b16d49b143853ab08b964(result, expected, **options):
    """Check if slide title matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '')
    expected_title = expected.get('expected_title', '')
    if actual_title is None:
        return 0.0
    if actual_title.strip().lower() == expected_title.strip().lower():
        return 1.0
    return 0.0

def check_slide_orientation__662bf4c59a75818f259c7e0aef709493(result, expected, **options):
    """Check if the slide orientation is Portrait (height > width)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    expected_orientation = expected.get('orientation', 'portrait')
    if expected_orientation == 'portrait':
        return 1.0 if height > width else 0.0
    else:
        return 1.0 if width > height else 0.0

def check_slide_title_font__05b860587b204cf23347a8d6b2a21320(result, expected, **options):
    """Check title text and font name on a slide. Partial credit: 0.5 text, 0.5 font."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    expected_font = expected.get('expected_font_name', '')
    actual_text = result.get('title_text')
    actual_font = result.get('font_name')
    if actual_text and expected_text:
        if actual_text.strip().lower() == expected_text.strip().lower():
            score += 0.5
    if actual_font and expected_font:
        if actual_font.strip().lower() == expected_font.strip().lower():
            score += 0.5
    return min(score, 1.0)

def check_slide_deleted__b3862aa49d8968cf4ed21032ba8b229f(result, expected, **options):
    """Check that a specific slide was deleted (count reduced, title gone)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 11)
    removed_title = expected.get('removed_title', 'Homework:')
    actual_count = result.get('slide_count', 0)
    titles = result.get('titles', [])
    if actual_count == expected_count:
        score += 0.5
    if removed_title not in titles:
        score += 0.5
    return min(score, 1.0)

def check_impress_new_slide__3a9a25804cb58719a25a23e52e7e124a(result, expected, **options):
    """Check if presentation has expected slide count and last slide has expected title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count')
    expected_last_title = expected.get('expected_last_title', '').lower()
    actual_count = result.get('slide_count', 0)
    slides = result.get('slides', [])
    if actual_count == expected_count:
        score += 0.5
    if slides and expected_last_title:
        last_title = (slides[-1].get('title') or '').lower()
        if expected_last_title in last_title:
            score += 0.5
    return min(score, 1.0)

def check_impress_table_with_content__608b93868302b656e50a31095c9d18e9(result, expected, **options):
    """Check table dimensions and first cell content."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not result.get('has_table'):
        return 0.0
    score = 0.0
    expected_rows = expected.get('expected_rows')
    expected_cols = expected.get('expected_cols')
    expected_text = expected.get('expected_first_cell', '').lower()
    if result.get('rows') == expected_rows:
        score += 0.4
    if result.get('cols') == expected_cols:
        score += 0.3
    actual_text = (result.get('first_cell_text') or '').lower()
    if expected_text and expected_text in actual_text:
        score += 0.3
    return min(score, 1.0)

def check_slide_orientation__326a96d7a18db487b955505da7fa674b(result, expected, **options):
    """Check if slide orientation matches expected (portrait or landscape)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_portrait = expected.get('is_portrait', True)
    actual_portrait = result.get('is_portrait', False)
    if actual_portrait == expected_portrait:
        return 1.0
    return 0.0

def check_slide_title__ebdab037a8eab6c57e56f5796ba7c0ab(result, expected, **options):
    """Check if the slide title matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title', '')
    if actual_title.strip().lower() == expected_title.strip().lower():
        return 1.0
    return 0.0

def check_pptx_pic_height_and_fonts__5b53387f975ab04f84e7061a6b89939a(result, expected, **options):
    """Check picture height and font sizes. Partial credit: 0.5 for each correct part."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_height = expected.get('expected_height_cm')
    expected_font_size = expected.get('expected_font_size_pt')
    tolerance = expected.get('tolerance_cm', 0.5)
    pic_height = result.get('pic_height_cm')
    if pic_height is not None and expected_height is not None:
        if abs(pic_height - expected_height) <= tolerance:
            score += 0.5
    font_sizes = result.get('font_sizes', [])
    if font_sizes and expected_font_size is not None:
        matching = sum((1 for s in font_sizes if abs(s - expected_font_size) <= 1.0))
        if len(font_sizes) > 0:
            ratio = matching / len(font_sizes)
            score += 0.5 * ratio
    return min(score, 1.0)

def check_impress_slide_deletion__1abdd5c6fe686f66f2606ee0a55bf0d1(result, expected, **options):
    """Check that slide deletion was performed correctly.

    Checks:
    - slide_count matches expected
    - first slide text contains expected substring (verifying correct slide was deleted)
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_count')
    expected_first_text_contains = expected.get('expected_first_text_contains', '')
    if result.get('slide_count') == expected_count:
        score += 0.5
    slide_texts = result.get('slide_texts', {})
    first_text = slide_texts.get('0', '')
    if first_text and expected_first_text_contains and (expected_first_text_contains.lower() in first_text.lower()):
        score += 0.5
    return min(score, 1.0)

def check_impress_title_and_bg__ac120f719f06d988c1c09c7022ee4b50(result, expected, **options):
    """Check slide 5 title and slide 6 background color.
    Partial credit: 0.5 for correct title, 0.5 for correct background.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_title = expected.get('expected_title', 'Game Rules')
    actual_title = result.get('slide5_title', '')
    if actual_title and expected_title.lower() in actual_title.lower():
        score += 0.5
    expected_color = expected.get('expected_bg_color', 'FF0000').upper()
    actual_color = (result.get('slide6_bg_color') or '').upper()
    if actual_color:
        try:
            er = int(expected_color[0:2], 16)
            eg = int(expected_color[2:4], 16)
            eb = int(expected_color[4:6], 16)
            ar = int(actual_color[0:2], 16)
            ag = int(actual_color[2:4], 16)
            ab = int(actual_color[4:6], 16)
            tolerance = expected.get('color_tolerance', 60)
            if abs(er - ar) <= tolerance and abs(eg - ag) <= tolerance and (abs(eb - ab) <= tolerance):
                score += 0.5
        except (ValueError, IndexError):
            pass
    return min(score, 1.0)

def check_pptx_image_size__19c7a30018144cafc9f4ebd4b22910ef(result, expected, **options):
    """Check if an image exists on the slide with expected dimensions."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    image_count = result.get('image_count', 0)
    if image_count == 0:
        return 0.0
    expected_width = expected.get('expected_width_cm')
    expected_height = expected.get('expected_height_cm')
    tolerance = expected.get('tolerance_cm', 0.05)
    score = 0.0
    score += 0.5
    for img in result.get('images', []):
        w_ok = abs(img['width_cm'] - expected_width) <= tolerance
        h_ok = abs(img['height_cm'] - expected_height) <= tolerance
        if w_ok and h_ok:
            score += 0.5
            break
    return min(score, 1.0)

def check_pptx_slide_count__f16a19b8aecc56a3153baf7aa9668d95(result, expected, **options):
    """Check if the slide count matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_count = result.get('slide_count')
    expected_count = expected.get('expected_count')
    if actual_count is None or expected_count is None:
        return 0.0
    return 1.0 if actual_count == expected_count else 0.0

def check_pptx_title_text__661ea64646d6a38833b5fdf5450470f8(result, expected, **options):
    """Check if the title text matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_text = result.get('title_text', '').strip()
    expected_text = expected.get('expected_text', '').strip()
    if actual_text.lower() == expected_text.lower():
        return 1.0
    return 0.0

def check_impress_all_title_fonts__9c7bd04930deae5cddd57fa6475996f4(result, expected, **options):
    """Check if all slide titles have the expected font. Partial credit per slide."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    fonts = result.get('fonts', {})
    slide_count = result.get('slide_count', 0)
    expected_font = expected.get('expected_font', 'Arial')
    if slide_count == 0:
        return 0.0
    correct = 0
    for (key, font) in fonts.items():
        if font == expected_font:
            correct += 1
    return correct / slide_count if slide_count > 0 else 0.0

def check_pptx_pic_height_and_fonts__4eb9ebb1f165db3d3eb60748b0824590(result, expected, **options):
    """Check picture height and font sizes. Partial credit: 0.5 for each correct part."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_height = expected.get('expected_height_cm')
    expected_font_size = expected.get('expected_font_size_pt')
    tolerance = expected.get('tolerance_cm', 0.5)
    pic_height = result.get('pic_height_cm')
    if pic_height is not None and expected_height is not None:
        if abs(pic_height - expected_height) <= tolerance:
            score += 0.5
    font_sizes = result.get('font_sizes', [])
    if font_sizes and expected_font_size is not None:
        matching = sum((1 for s in font_sizes if abs(s - expected_font_size) <= 1.0))
        if len(font_sizes) > 0:
            ratio = matching / len(font_sizes)
            score += 0.5 * ratio
    return min(score, 1.0)

def check_snapshot_and_slide_bg__da1a5547fec970e705ee94d1f2bd1e91(result, expected, **options):
    """Check snapshot exists and slide background was set. Partial credit."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('snapshot_exists', False):
        score += 0.4
    if result.get('has_bg_image', False):
        score += 0.6
    return min(score, 1.0)

def check_impress_strikethrough__941b469f6d56a3f2848b88f32a6979cb(result, expected, **options):
    """Check if target paragraphs have strikethrough formatting."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    target_texts = expected.get('target_texts', [])
    if not target_texts:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(target_texts)
    for target in target_texts:
        for para in paragraphs:
            if target.lower() in para['text'].lower():
                if para['strikethrough']:
                    score += per_item
                break
    return min(score, 1.0)

def check_pptx_orientation__f921e5b58a724c7baa4b59415c1c4019(result, expected, **options):
    """Check if the slide orientation matches expected (portrait or landscape)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_orientation = expected.get('orientation', 'portrait')
    is_portrait = result.get('is_portrait', False)
    if expected_orientation == 'portrait':
        return 1.0 if is_portrait else 0.0
    else:
        return 1.0 if not is_portrait else 0.0

def check_slide_title_font_color__66374a642830d11922d74321b71d53a4(result, expected, **options):
    """Check if the title font color matches the expected color within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_color = result.get('color', [])
    expected_color = expected.get('expected_color', [])
    if not actual_color or not expected_color:
        return 0.0
    if len(actual_color) != 3 or len(expected_color) != 3:
        return 0.0
    dist = sum(((a - e) ** 2 for (a, e) in zip(actual_color, expected_color))) ** 0.5
    threshold = expected.get('threshold', 30)
    if dist <= threshold:
        return 1.0
    return 0.0

def check_slide_dimensions__54270a3c8d106e03ce2d0c72089a3417(result, expected, **options):
    """Check if slide dimensions match expected width and height in cm.

    Partial credit: 0.5 for correct width, 0.5 for correct height.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance_cm', 0.5)
    score = 0.0
    expected_width = expected.get('expected_width_cm')
    expected_height = expected.get('expected_height_cm')
    actual_width = result.get('width_cm', 0)
    actual_height = result.get('height_cm', 0)
    if expected_width is not None and abs(actual_width - expected_width) <= tolerance:
        score += 0.5
    if expected_height is not None and abs(actual_height - expected_height) <= tolerance:
        score += 0.5
    return score

def check_pptx_slide_title__24e166a4a3628d77570c2919341fa3d0(result, expected, **options):
    """Check if the slide title matches expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title_text', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if actual_title == expected_title:
        return 1.0
    if actual_title.lower() == expected_title.lower():
        return 0.8
    return 0.0

def check_presentation_summary__f7a75c940bb0ba7687c094e66be96f16(result, expected, **options):
    """Check that presentation summary contains required information.

    Checks:
    - Total slide count is mentioned
    - Number of slides with notes is mentioned
    - Slide numbers that have notes are listed
    Partial credit for each component.
    """
    if result.get('error'):
        return 0.0
    content = result.get('content', '').lower()
    if not content:
        return 0.0
    score = 0.0
    total_checks = 3
    expected_total = str(expected.get('total_slides', 0))
    expected_notes_count = str(expected.get('notes_count', 0))
    expected_note_slides = expected.get('slides_with_notes', [])
    if expected_total in content:
        score += 1.0 / total_checks
    if expected_notes_count in content:
        score += 1.0 / total_checks
    if expected_note_slides:
        found = 0
        for slide_num in expected_note_slides:
            s = str(slide_num)
            if s in content:
                found += 1
        fraction = found / len(expected_note_slides)
        if fraction >= 0.5:
            score += 1.0 / total_checks
    return min(score, 1.0)

def check_impress_title_fontsize__f67d7e69d6a85882636e0a41dabd5c89(result, expected, **options):
    """Check if the slide title has the expected font size."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_pt = result.get('font_size_pt')
    expected_pt = expected.get('expected_size_pt', 36)
    if actual_pt is None:
        return 0.0
    if abs(actual_pt - expected_pt) < 0.5:
        return 1.0
    return 0.0

def check_pptx_slide_count_and_title__e6d2505b7a2f1323c83204293d9b2ff3(result, expected, **options):
    """Check if slide count and last slide title match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 3)
    actual_count = result.get('slide_count', 0)
    if actual_count == expected_count:
        score += 0.5
    expected_title = expected.get('expected_last_title', '').strip().lower()
    actual_title = result.get('last_slide_title', '').strip().lower()
    if expected_title and actual_title == expected_title:
        score += 0.5
    return min(score, 1.0)

def check_slide_title_text__f01c4fcc320219e33dd1b3dfc752a5b6(result, expected, **options):
    """Check if slide title text matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if not expected_title:
        return 0.0
    if actual_title.lower() == expected_title.lower():
        return 1.0
    return 0.0

def check_slide_count__55d37620fcf7998461b41e2b1579f16b(result, expected, **options):
    """Check if the slide count matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_count = expected.get('expected_count', 0)
    actual_count = result.get('slide_count', 0)
    if actual_count == expected_count:
        return 1.0
    return 0.0

def check_slide_notes__44fd5b3a7c457950926c4afa82f02524(result, expected, **options):
    """Check if slide notes contain the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    notes_text = result.get('notes_text', '')
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    if expected_text.lower() in notes_text.lower():
        return 1.0
    return 0.0

def check_slide2_underline_48pt__203dc282bb468063e328a7afef0f556d(result, expected, **options):
    """Check that welcome message on slide 2 is underlined and 48pt."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    runs = result.get('runs', [])
    if not runs:
        return 0.0
    score = 0.0
    expected_size = expected.get('expected_size_pt', 48.0)
    all_underline = all((r.get('underline') is True for r in runs))
    if all_underline:
        score += 0.5
    all_correct_size = all((r.get('size_pt') is not None and abs(r['size_pt'] - expected_size) < 0.5 for r in runs))
    if all_correct_size:
        score += 0.5
    return score

def check_pptx_pic_width_and_fonts__1a124508c5c52deecd8c9cfabb1cdd16(result, expected, **options):
    """Check picture width and font sizes. Partial credit: 0.5 for each correct part."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_width = expected.get('expected_width_cm')
    expected_font_size = expected.get('expected_font_size_pt')
    tolerance = expected.get('tolerance_cm', 0.5)
    pic_width = result.get('pic_width_cm')
    if pic_width is not None and expected_width is not None:
        if abs(pic_width - expected_width) <= tolerance:
            score += 0.5
    font_sizes = result.get('font_sizes', [])
    if font_sizes and expected_font_size is not None:
        matching = sum((1 for s in font_sizes if abs(s - expected_font_size) <= 1.0))
        if len(font_sizes) > 0:
            ratio = matching / len(font_sizes)
            score += 0.5 * ratio
    return min(score, 1.0)

def check_slide_text__64fdb0e52d9e147a698cbe6791191c81(result, expected, **options):
    """Check if a specific text exists on the slide (and optionally an old text is gone)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    texts = result.get('texts', [])
    expected_text = expected.get('expected_text', '')
    old_text = expected.get('old_text', '')
    score = 0.0
    for t in texts:
        if expected_text.lower() in t.lower():
            score += 0.5
            break
    if old_text:
        old_found = False
        for t in texts:
            if old_text.lower() in t.lower():
                old_found = True
                break
        if not old_found:
            score += 0.5
    else:
        score += 0.5
    return min(score, 1.0)

def check_slide_title__7b96cb41298a29a8e5448750424ab134(result, expected, **options):
    """Check if the slide title matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '')
    expected_title = expected.get('expected_title', '')
    if not actual_title or not expected_title:
        return 0.0
    if actual_title.strip().lower() == expected_title.strip().lower():
        return 1.0
    if expected_title.strip().lower() in actual_title.strip().lower():
        return 0.5
    return 0.0

def check_slide_title__94c866527f562343eb1ed6d961346913(result, expected, **options):
    """Check if slide title text and alignment match expected values.
    Partial credit: 0.5 for correct text, 0.5 for correct alignment.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    expected_alignment = expected.get('expected_alignment', '')
    actual_text = result.get('text', '').strip()
    actual_alignment = result.get('alignment', '')
    if actual_text.lower() == expected_text.lower():
        score += 0.5
    if actual_alignment == expected_alignment:
        score += 0.5
    return score

def check_slide5_bold_36pt__e80588656c5c992160d4d54a69813615(result, expected, **options):
    """Check that title on slide 5 is bold and 36pt."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    runs = result.get('runs', [])
    if not runs:
        return 0.0
    score = 0.0
    expected_size = expected.get('expected_size_pt', 36.0)
    all_bold = all((r.get('bold') is True for r in runs))
    if all_bold:
        score += 0.5
    all_correct_size = all((r.get('size_pt') is not None and abs(r['size_pt'] - expected_size) < 0.5 for r in runs))
    if all_correct_size:
        score += 0.5
    return score

def check_slide_bg_color__147ec12b6038f2d22f20700395a41639(result, expected, **options):
    """Check if the slide background color matches the expected RGB value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_color = result.get('color_rgb', '').upper()
    expected_color = expected.get('expected_color', '').upper()
    if not actual_color or not expected_color or actual_color == 'NONE':
        return 0.0
    if actual_color == expected_color:
        return 1.0
    try:
        (ar, ag, ab) = (int(actual_color[0:2], 16), int(actual_color[2:4], 16), int(actual_color[4:6], 16))
        (er, eg, eb) = (int(expected_color[0:2], 16), int(expected_color[2:4], 16), int(expected_color[4:6], 16))
        distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
        if distance < 30:
            return 0.5
    except (ValueError, IndexError):
        pass
    return 0.0

def check_slide_title__497f89854667fa402a3b0a0e21fc061d(result, expected, **options):
    """Check if slide title text and alignment match expected values.
    Partial credit: 0.5 for correct text, 0.5 for correct alignment.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_text', '')
    expected_alignment = expected.get('expected_alignment', '')
    actual_text = result.get('text', '').strip()
    actual_alignment = result.get('alignment', '')
    if actual_text.lower() == expected_text.lower():
        score += 0.5
    if actual_alignment == expected_alignment:
        score += 0.5
    return score

def check_slide_title__07d639643e02063bb060c675d5eb936b(result, expected, **options):
    """Check if the slide title matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_title = result.get('title', '')
    expected_title = expected.get('expected_title', '')
    if actual_title is None:
        return 0.0
    if actual_title.strip().lower() == expected_title.strip().lower():
        return 1.0
    return 0.0

def check_impress_slide_duplicate__b84f4f9ed96e0e954dfada4f0a280734(result, expected, **options):
    """Check that slide duplication was performed correctly.

    Checks:
    - slide_count matches expected
    - last slide text contains expected substring (matching duplicated slide)
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_count')
    expected_last_text_contains = expected.get('expected_last_text_contains', '')
    if result.get('slide_count') == expected_count:
        score += 0.5
    slide_texts = result.get('slide_texts', {})
    last_text = slide_texts.get('-1', '')
    if last_text and expected_last_text_contains and (expected_last_text_contains.lower() in last_text.lower()):
        score += 0.5
    return min(score, 1.0)

def check_slide_titles__d445ea0218043b2b599f555fb0b0c923(result, expected, **options):
    """Check that extracted slide titles match expected titles."""
    if result.get('error'):
        return 0.0
    actual_lines = result.get('lines', [])
    expected_titles = expected.get('expected_titles', [])
    if not expected_titles:
        return 0.0
    total = len(expected_titles)
    matched = 0
    for (i, exp_title) in enumerate(expected_titles):
        if i < len(actual_lines):
            if actual_lines[i].strip().lower() == exp_title.strip().lower():
                matched += 1
    return matched / total if total > 0 else 0.0

def check_pptx_title_bold__7728836b00929340fdef9a2a6602ef09(result, expected, **options):
    """Check if the title text is bold."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_bold = expected.get('expected_bold', True)
    actual_bold = result.get('all_bold', False)
    if actual_bold == expected_bold:
        return 1.0
    return 0.0

def check_impress_slide_table__dabff5cb0982227a6c01462d5c667c85(result, expected, **options):
    """Check if a slide has a table with expected rows and columns."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tables = result.get('tables', [])
    if not tables:
        return 0.0
    expected_rows = expected.get('expected_rows')
    expected_cols = expected.get('expected_cols')
    for tbl in tables:
        if tbl.get('rows') == expected_rows and tbl.get('cols') == expected_cols:
            return 1.0
    if tables:
        return 0.3
    return 0.0

def check_pptx_paragraph_contains__10c4d6d7b532dd358d2936a84f08db4a(result, expected, **options):
    """Check if a specific text exists among the slide paragraphs."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    required_text = expected.get('required_text', '')
    if not required_text:
        return 0.0
    for para in paragraphs:
        if required_text.lower() in para.lower():
            return 1.0
    return 0.0

def check_pptx_image_size__32c4d5ced54cde2e00ee07b6ba3973e4(result, expected, **options):
    """Check if an image exists on the slide with expected dimensions."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    image_count = result.get('image_count', 0)
    if image_count == 0:
        return 0.0
    expected_width = expected.get('expected_width_cm')
    expected_height = expected.get('expected_height_cm')
    tolerance = expected.get('tolerance_cm', 0.05)
    score = 0.0
    score += 0.5
    for img in result.get('images', []):
        w_ok = abs(img['width_cm'] - expected_width) <= tolerance
        h_ok = abs(img['height_cm'] - expected_height) <= tolerance
        if w_ok and h_ok:
            score += 0.5
            break
    return min(score, 1.0)

def check_pptx_all_slides_bg_color__88a6aaaf9999692a07033829e3809d82(result, expected, **options):
    """Check if all slides have the expected background color. Partial credit per slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    slides = result.get('slides', [])
    total = result.get('total_slides', 0)
    if total == 0:
        return 0.0
    expected_color = expected.get('expected_color', '').upper()
    tolerance = expected.get('tolerance', 30)
    matching = 0
    for slide_info in slides:
        color_rgb = slide_info.get('color_rgb')
        if color_rgb is None:
            continue
        actual_color = color_rgb.upper()
        if actual_color == expected_color:
            matching += 1
            continue
        try:
            ar = int(actual_color[0:2], 16)
            ag = int(actual_color[2:4], 16)
            ab = int(actual_color[4:6], 16)
            er = int(expected_color[0:2], 16)
            eg = int(expected_color[2:4], 16)
            eb = int(expected_color[4:6], 16)
            distance = ((ar - er) ** 2 + (ag - eg) ** 2 + (ab - eb) ** 2) ** 0.5
            if distance <= tolerance:
                matching += 1
        except (ValueError, IndexError):
            pass
    return matching / total

def check_pptx_textbox_font_sizes__d842f0d6650bc2853fbd82b34d9bb662(result, expected, **options):
    """Check if textbox font sizes match expected values. Supports partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    font_sizes = result.get('font_sizes', {})
    expected_sizes = expected.get('expected_sizes', {})
    if not expected_sizes:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_sizes)
    for (text, expected_size) in expected_sizes.items():
        actual_size = font_sizes.get(text)
        if actual_size is not None and abs(actual_size - expected_size) < 0.5:
            score += per_item
    return min(score, 1.0)

def check_impress_title_font_props__8c5c9bdd88c0ba0d1f43df42a252cfae(result, expected, **options):
    """Check if title has expected italic and font size properties.

    Partial credit: 0.5 for italic, 0.5 for size.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_italic = expected.get('expected_italic', True)
    actual_italic = result.get('italic')
    if actual_italic == expected_italic:
        score += 0.5
    expected_size = expected.get('expected_size_pt')
    actual_size = result.get('size_pt')
    if expected_size is not None and actual_size is not None:
        if abs(float(actual_size) - float(expected_size)) <= 1.0:
            score += 0.5
    return min(score, 1.0)

def check_pptx_shape_dimensions__37468b49cfdc53eaca3dfbbee86985a3(result, expected, **options):
    """Check if shape dimensions match expected values with tolerance and partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    tolerance = expected.get('tolerance', 0.5)
    score = 0.0
    weight = 1.0 / len(checks)
    for check in checks:
        key = check['key']
        expected_value = check['value']
        actual = result.get(key)
        if actual is not None and abs(actual - expected_value) <= tolerance:
            score += weight
    return min(round(score, 4), 1.0)

def check_impress_bg_and_subtitle__024a68353adff6dd2d6a01cc842a46ed(result, expected, **options):
    """Check slides 1&6 background color and slide 2 subtitle text.
    Partial credit: 0.34 for slide 1 bg, 0.33 for slide 6 bg, 0.33 for subtitle.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_color = expected.get('expected_bg_color', 'FFA500').upper()
    tolerance = expected.get('color_tolerance', 60)
    slide1_hex = result.get('slide1_bg_color')
    if slide1_hex:
        slide1_hex = slide1_hex.upper()
        try:
            (er, eg, eb) = (int(expected_color[0:2], 16), int(expected_color[2:4], 16), int(expected_color[4:6], 16))
            (ar, ag, ab) = (int(slide1_hex[0:2], 16), int(slide1_hex[2:4], 16), int(slide1_hex[4:6], 16))
            if abs(er - ar) <= tolerance and abs(eg - ag) <= tolerance and (abs(eb - ab) <= tolerance):
                score += 0.34
        except (ValueError, IndexError):
            pass
    slide6_hex = result.get('slide6_bg_color')
    if slide6_hex:
        slide6_hex = slide6_hex.upper()
        try:
            (er, eg, eb) = (int(expected_color[0:2], 16), int(expected_color[2:4], 16), int(expected_color[4:6], 16))
            (ar, ag, ab) = (int(slide6_hex[0:2], 16), int(slide6_hex[2:4], 16), int(slide6_hex[4:6], 16))
            if abs(er - ar) <= tolerance and abs(eg - ag) <= tolerance and (abs(eb - ab) <= tolerance):
                score += 0.33
        except (ValueError, IndexError):
            pass
    expected_subtitle = expected.get('expected_subtitle', 'Welcome to our presentation')
    actual_subtitle = result.get('slide2_subtitle', '')
    if actual_subtitle and expected_subtitle.lower() in actual_subtitle.lower():
        score += 0.33
    return min(score, 1.0)

def check_pptx_transitions__a8621b66c753635f21d8c612a865d275(result, expected, **options):
    """Check if all slides have the expected transition type. Returns partial credit."""
    if isinstance(result, str) or not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_type = expected.get('transition_type', 'fade').lower()
    transitions = result.get('transitions', [])
    total = result.get('total_slides', 0)
    if total == 0:
        return 0.0
    correct = sum((1 for t in transitions if t is not None and t == expected_type))
    return correct / total

def check_slide_image_inserted__e8973b7aba53f626fa037406b1da4c8e(result, expected, **options):
    """Check that the first slide has an image and no text placeholders.
    Partial credit: 0.5 for image present, 0.5 for no text placeholders.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    min_images = expected.get('min_images', 1)
    image_count = result.get('image_count', 0)
    if image_count >= min_images:
        score += 0.5
    text_placeholder_count = result.get('text_placeholder_count', -1)
    if text_placeholder_count == 0:
        score += 0.5
    return min(score, 1.0)

def check_impress_title_bold__6e02436ff918c1445d486be48242a3ca(result, expected, **options):
    """Check if the slide title has the expected bold property."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual_bold = result.get('bold', False)
    expected_bold = expected.get('expected_bold', True)
    return 1.0 if actual_bold == expected_bold else 0.0

def check_pptx_multi_slide_images__e6e9ce905cbd4fd74addd1fc2cd640de(result, expected, **options):
    """Check if images exist on multiple slides with expected dimensions.
    Partial credit: 0.5 per slide (2 slides = 1.0).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    slides_data = result.get('slides', {})
    expected_width = expected.get('expected_width_cm')
    expected_height = expected.get('expected_height_cm')
    tolerance = expected.get('tolerance_cm', 0.05)
    slide_indices = expected.get('slide_indices', ['0', '1'])
    score = 0.0
    per_slide_score = 1.0 / len(slide_indices)
    for idx_str in slide_indices:
        slide_data = slides_data.get(idx_str, {})
        if slide_data.get('error') or slide_data.get('image_count', 0) == 0:
            continue
        for img in slide_data.get('images', []):
            w_ok = abs(img['width_cm'] - expected_width) <= tolerance
            h_ok = abs(img['height_cm'] - expected_height) <= tolerance
            if w_ok and h_ok:
                score += per_slide_score
                break
    return min(score, 1.0)

def check_slide_bg_color__1dbf2adcf932ab2d0e517ac3d78f9d08(result, expected, **options):
    """Check if slide background color matches expected RGB within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict) or not result.get('has_color'):
        return 0.0
    expected_r = expected.get('r', 0)
    expected_g = expected.get('g', 0)
    expected_b = expected.get('b', 0)
    threshold = expected.get('threshold', 50)
    actual_r = result.get('r', 255)
    actual_g = result.get('g', 255)
    actual_b = result.get('b', 255)
    distance = ((actual_r - expected_r) ** 2 + (actual_g - expected_g) ** 2 + (actual_b - expected_b) ** 2) ** 0.5
    if distance <= threshold:
        return 1.0
    elif distance <= threshold * 2:
        return 0.5
    return 0.0

def check_audio_on_slide__5688a332db768af3543ab15e11555c76(result, expected, **options):
    """Check if audio is present on the specified slide.

    Result (from getter):
        has_audio: bool
        audio_count: int
        slide_count: int

    Expected (from rules):
        min_audio_count: minimum number of audio files expected (default 1)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    min_count = expected.get('min_audio_count', 1)
    audio_count = result.get('audio_count', 0)
    if audio_count >= min_count:
        return 1.0
    return 0.0

def check_slide_count__3e661325ce185501287a57bdd6c09b52(result, expected, **options):
    """Check slide count after deletion and file existence."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error') or not result.get('exists'):
        return 0.0
    score = 0.0
    if result.get('exists'):
        score += 0.3
    expected_count = expected.get('expected_slide_count', 15)
    actual_count = result.get('slide_count', 0)
    if actual_count == expected_count:
        score += 0.7
    elif abs(actual_count - expected_count) <= 1:
        score += 0.35
    return min(score, 1.0)

def check_slide_count_blank__c9d9670d65b42979a3ba7969e086139a(result, expected, **options):
    """Check slide count and that all slides are blank (no text shapes).
    Partial credit: 0.5 for correct count, 0.5 for all blank.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 6)
    actual_count = result.get('slide_count', 0)
    if actual_count == expected_count:
        score += 0.5
    slides_with_text = result.get('slides_with_text', -1)
    if slides_with_text == 0:
        score += 0.5
    return min(score, 1.0)

def check_pptx_italic__979792a60bf386d8e1c5a4d702d19829(result, expected, **options):
    """Check if text on specified slides has italic formatting. Partial credit per slide."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    slide_indices = expected.get('slide_indices', [0, 1])
    score = 0.0
    weight = 1.0 / len(slide_indices) if slide_indices else 0.0
    for idx in slide_indices:
        key = f'slide_{idx}'
        slide_data = result.get(key, {})
        if slide_data.get('error'):
            continue
        total = slide_data.get('total_runs', 0)
        italic = slide_data.get('italic_runs', 0)
        if total > 0 and italic == total:
            score += weight
    return min(score, 1.0)

def check_impress_slide_reorder__88930a069947104f21bb2b61fdebcc2c(result, expected, **options):
    """Check that slide reordering was performed correctly.

    Checks:
    - slide_count matches expected
    - specific slide positions contain expected text substrings
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_count')
    text_checks = expected.get('text_checks', {})
    if result.get('slide_count') == expected_count:
        score += 0.34
    slide_texts = result.get('slide_texts', {})
    num_checks = len(text_checks)
    if num_checks > 0:
        points_per_check = 0.66 / num_checks
        for (idx_str, expected_substr) in text_checks.items():
            actual_text = slide_texts.get(idx_str, '')
            if actual_text and expected_substr.lower() in actual_text.lower():
                score += points_per_check
    return min(score, 1.0)

def check_pptx_title_text__927f9acbea3f73fbaa9e281542feb428(result, expected, **options):
    """Check if the title text matches the expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_text = expected.get('expected_text', '')
    actual_text = result.get('title_text', '')
    if actual_text.strip().lower() == expected_text.strip().lower():
        return 1.0
    return 0.0

def check_slide_title_and_color__fd2e510c274f3f89d65383d731b18828(result, expected, **options):
    """Check slide title text and font color with partial credit.

    Scoring: 0.5 for correct title text, 0.5 for correct font color.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_title = result.get('title', '').strip()
    expected_title = expected.get('expected_title', '').strip()
    if expected_title and actual_title.lower() == expected_title.lower():
        score += 0.5
    actual_color = result.get('font_color', '')
    expected_color = expected.get('expected_color', '')
    if expected_color and actual_color:
        if actual_color.upper() == expected_color.upper():
            score += 0.5
    return min(score, 1.0)

def check_slide_count__1573725569136f71d4ad79ac5ffbf35a(result, expected, **options):
    """Check if slide count matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_count = result.get('slide_count', -1)
    expected_count = expected.get('expected_count', -1)
    if actual_count == expected_count:
        return 1.0
    return 0.0

def check_slide_duplicated__e9203909bdddcf8283496218733a0b9b(result, expected, **options):
    """Check that slide 1 was duplicated (total=13, slide 2 title matches slide 1)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 13)
    expected_title = expected.get('expected_duplicate_title', 'Forests')
    actual_count = result.get('slide_count', 0)
    slide2_title = result.get('slide2_title', '')
    if actual_count == expected_count:
        score += 0.5
    if slide2_title == expected_title:
        score += 0.5
    return min(score, 1.0)

def check_pptx_slide_count__d96a0661e83ae9239b523060c49e36f2_qw35sft2_4679f70e(result, expected, **options):
    """Check that the PPTX has the expected number of slides."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_count = expected.get('expected_count', 2)
    actual_count = result.get('slide_count', 0)
    return 1.0 if actual_count == expected_count else 0.0

def check_pptx_table_row0__27882b2f85f9d8345d5e8153a8f07379_qw35sft2_bf84994b(result, expected, **options):
    """Check if the first row of the table matches expected values. All 4 cells must match."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_row = expected.get('expected_table_row0', [])
    actual_row = result.get('table_row0', [])
    if not expected_row or not actual_row:
        return 0.0
    if len(actual_row) != len(expected_row):
        return 0.0
    matches = sum((1 for a, e in zip(actual_row, expected_row) if a == e))
    return matches / len(expected_row)

def check_pptx_summary_notes__f1eef050d638a5408c254c4e620a5302_qw35sft2_a99f5872(result, expected, **options):
    """Partial credit: 0.5 for correct slide count, 0.5 for expected keyword in last slide notes."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 13)
    if result.get('slide_count') == expected_count:
        score += 0.5
    keyword = expected.get('expected_notes_keyword', '').lower()
    actual_notes = result.get('last_slide_notes', '').lower()
    if keyword and keyword in actual_notes:
        score += 0.5
    return min(score, 1.0)

def check_pptx_last_moved_to_first__b9dba6723fe45eeec63ed20d4f46d22d_qw35sft2_376a34c4(result, expected, **options):
    """Check that the last slide was moved to the first position."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_count = expected.get('expected_slide_count', 24)
    expected_keyword = expected.get('first_slide_keyword', 'Works Cited')
    actual_count = result.get('slide_count', 0)
    first_text = result.get('first_slide_text', '')
    if actual_count != expected_count:
        return 0.0
    if expected_keyword.lower() not in first_text.lower():
        return 0.0
    return 1.0

def check_impress_title_bottom_and_text__dea2c3a0c75be4ec786c2a7bbd9ea59f_qw35sft2_ff64be5a(result, expected, **options):
    """Partial credit: 0.5 for title at bottom half, 0.5 for title text matching expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    top_cm = result.get('top_cm')
    min_top_cm = expected.get('min_top_cm', 9.525)
    if top_cm is not None and top_cm >= min_top_cm:
        score += 0.5
    title_text = result.get('title_text', '')
    expected_text = expected.get('expected_text', '')
    if expected_text and title_text.strip() == expected_text.strip():
        score += 0.5
    return min(score, 1.0)

def check_pptx_slide_transitions__f348aaf139c455284fcbda7cfe577da5_qw35sft2_a8a2fc6d(result, expected, **options):
    """Check if all slides have the expected transition type applied."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_type = expected.get('transition_type', 'fade').lower()
    transitions = result.get('transitions', [])
    total = result.get('total_slides', 0)
    if total == 0:
        return 0.0
    matching = sum((1 for t in transitions if t.lower() == expected_type))
    return matching / total

def check_pptx_image_size__7a22b4a6c189fa88dd4ad54ac62a04cd_qw35sft2_5933753b(result, expected, **options):
    """Check if a slide has at least one image matching the expected size (with tolerance)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    images = result.get('images', [])
    min_count = expected.get('min_image_count', 1)
    if len(images) < min_count:
        return 0.0
    exp_w = float(expected.get('expected_width_cm', 1.0))
    exp_h = float(expected.get('expected_height_cm', 1.0))
    tol = float(expected.get('tolerance_cm', 0.15))
    for img in images:
        w = float(img.get('width_cm', 0))
        h = float(img.get('height_cm', 0))
        if abs(w - exp_w) <= tol and abs(h - exp_h) <= tol:
            return 1.0
    return 0.0

def check_pptx_bg_solid__7f279c8dbe68c5780c5b5f60310b951b_qw35sft2_a43b9793(result, expected, **options):
    """Check if slide background has been changed to a solid color fill (SOLID type).

    The initial file has fill type BACKGROUND (master/theme). Any solid color
    assignment will yield SOLID, satisfying this check.
    """
    if result.get('error'):
        return 0.0
    fill_type = result.get('fill_type', '')
    expected_fill = expected.get('expected_fill_type', 'SOLID')
    return 1.0 if fill_type == expected_fill else 0.0

def check_all_slides_fade__85331dc0cf4d821fa2bad3e27657d683_qw35sft2_3c31d6ea(result, expected, **options):
    """Check that all slides in the presentation have a Fade transition.

    LibreOffice Impress exports a fade as <p:fade> inside <p:transition>.
    The getter returns type='fade' for each slide if this element is present.

    Scoring:
      - 1.0  if ALL slides have a fade transition
      - partial credit (fraction of slides) otherwise

    expected keys:
      - transition_type: expected transition element name (default 'fade')
      - min_fraction: minimum fraction required for full credit (default 1.0)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    transitions = result.get('transitions', [])
    num_slides = result.get('num_slides', 0)
    if num_slides == 0 or not transitions:
        return 0.0
    expected_type = expected.get('transition_type', 'fade').lower()
    min_fraction = expected.get('min_fraction', 1.0)
    count = sum((1 for t in transitions if t.get('type', '').lower() == expected_type))
    fraction = count / num_slides
    if fraction >= min_fraction:
        return 1.0
    return round(fraction, 4)

def check_pptx_slide_bg_color__1a5b2a2b44ee739e3486ddbc20ac1c87_qw35sft2_36618094(result, expected, **options):
    """Check if the slide background color is approximately yellow."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if result.get('color_hex') is None:
        return 0.0
    expected_hex = expected.get('expected_color', 'FFFF00').upper().lstrip('#')
    tolerance = expected.get('tolerance', 30)
    try:
        exp_r = int(expected_hex[0:2], 16)
        exp_g = int(expected_hex[2:4], 16)
        exp_b = int(expected_hex[4:6], 16)
        act_r = result.get('r', 0)
        act_g = result.get('g', 0)
        act_b = result.get('b', 0)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_pptx_slide_has_solid_bg__30ec500dc9631258b0ce307c50b98145_qw35sft2_6aa292db(result, expected, **options):
    """Check if the slide has a solid color background fill (any color)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('is_solid') else 0.0

def check_impress_slide4_contact_shapes__5eae632a29df23abd3f96eee0032e6dc_qw35sft2_3d109fc7(result, expected, **options):
    """
    Check that the phone number has been removed from slide 4.
    Returns 1.0 if phone is absent, 0.0 otherwise.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('has_phone', True):
        return 1.0
    return 0.0

def check_impress_underline__fa317ec6034191f7b02a906f78ff4b09_qw35sft2_7c170663(result, expected, **options):
    """Check that all text runs in the target shape are underlined."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if result.get('all_underlined') is True:
        return 1.0
    return 0.0

def check_impress_audio_slide1_trans_slide3__1b01b098ee4b6dfab8bc1b2d24723a1f_qw35sft2_f2414360(result, expected, **options):
    """Partial credit: 0.5 for audio in slide 1, 0.5 for transition on slide 3."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('has_audio_slide1'):
        score += 0.5
    if result.get('has_transition_slide3'):
        score += 0.5
    return min(score, 1.0)

def check_blank_slide_count__da57195bec8b00f5d49d80dfe28b0049_qw35sft2_946e4f58(result, expected, **options):
    """Check that presentation has expected slide count and all slides are blank (no text frames)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('slide_count')
    if expected_count is not None and result.get('slide_count') == expected_count:
        score += 0.5
    if result.get('slides_with_text_frames', 999) == 0:
        score += 0.5
    return score

def check_impress_bullet_newpara__12c621423b2418d27e86fd746731530a_qw35sft2_85f3aeeb(result, expected, **options):
    """Check bullet on first para + second para contains expected text: 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('first_para_has_bullet') is True:
        score += 0.5
    expected_text = expected.get('expected_second_para_text', '').strip().lower()
    actual_text = result.get('second_para_text', '').strip().lower()
    if expected_text and expected_text in actual_text:
        score += 0.5
    return score

def check_pptx_text_alignments__6253f908caa7dec8dc9ccbad6c2f7332_qw35sft2_756004b9(result, expected, **options):
    """Check text alignment for slides 3, 4, 5. Partial credit: 1/3 per slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    per_slide = 1.0 / 3
    if result.get('slide3_align') == expected.get('slide3_align'):
        score += per_slide
    if result.get('slide4_align') == expected.get('slide4_align'):
        score += per_slide
    if result.get('slide5_align') == expected.get('slide5_align'):
        score += per_slide
    return min(round(score, 4), 1.0)

def check_pptx_slide_font_size__4b92c85bf54981a99e31cdf8daa555a4_qw35sft2_3d7ddbc4(result, expected, **options):
    """Check font name (0.5) and font size in pt (0.5) of last slide title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('font_name') == expected.get('font_name'):
        score += 0.5
    expected_size = expected.get('font_size_pt')
    actual_size = result.get('font_size_pt')
    if expected_size is not None and actual_size is not None:
        if abs(actual_size - expected_size) <= 1:
            score += 0.5
    return score

def check_pptx_slide3_text_colors__9e2b961f79a8b89da2fe1a7f40c57f65_qw35sft2_0a01b8df(result, expected, **options):
    """Check that all text runs in slide 3 have the expected font color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    colors = result.get('colors', [])
    if not colors:
        return 0.0
    expected_color = expected.get('expected_color', 'FFFF00').upper().lstrip('#')
    all_match = all((c == expected_color for c in colors if c is not None))
    none_present = any((c is None for c in colors))
    if all_match and (not none_present):
        return 1.0
    return 0.0

def check_pptx_slide3_para_texts__c0cb167c6059fd8f9f81e328b9599144_qw35sft2_a0cf6b20(result, expected, **options):
    """Check that target paragraph text was changed to expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    paragraphs = result.get('paragraphs', [])
    expected_text = expected.get('expected_text', '')
    old_text = expected.get('old_text', '')
    has_new = any((t == expected_text for t in paragraphs))
    no_old = all((t != old_text for t in paragraphs))
    if has_new and no_old:
        return 1.0
    if has_new:
        return 0.7
    return 0.0

def check_impress_multi_title_color__14718ed4f8dc81749072a781b972112d_qw35sft2_e6acc344(result, expected, **options):
    """
    Check if titles in multiple slides all have the expected color.
    Partial credit: 1/N per correctly colored slide.
    expected keys: expected_color_hex, slide_indices (list of 0-based ints)
    Returns: 0.0 to 1.0
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    slides_data = result.get('slides', {})
    if not slides_data:
        return 0.0
    expected_color = str(expected.get('expected_color_hex', '000000')).upper().lstrip('#')
    slide_indices = expected.get('slide_indices', list(slides_data.keys()))
    total = len(slide_indices)
    if total == 0:
        return 0.0
    correct = 0
    for idx in slide_indices:
        slide_info = slides_data.get(str(idx), {})
        if slide_info.get('error'):
            continue
        all_colors = slide_info.get('all_colors', [slide_info.get('color_hex')])
        if all_colors and all((c is not None and c.upper().lstrip('#') == expected_color for c in all_colors)):
            correct += 1
    return correct / total

def check_impress_slide1_title__992b8b598a921e5736a852e261be2237_qw35sft2_e1bfdb4f(result, expected, **options):
    """Check if slide 1 title text matches the expected value."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_text = expected.get('expected_title', '')
    actual_text = result.get('title_text', '')
    if not expected_text:
        return 0.0
    if actual_text.strip() == expected_text.strip():
        return 1.0
    return 0.0

def check_pptx_slide2_title_state__1225c3617161e6b819a6b6275a5c9749_qw35sft2_d3d81117(result, expected, **options):
    """Check slide 2 title text is 'Note' and alignment is RIGHT."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('title_text', 'Note')
    expected_alignment = expected.get('title_alignment', 'RIGHT')
    actual_text = result.get('title_text', '')
    actual_alignment = result.get('title_alignment', '')
    if actual_text == expected_text:
        score += 0.5
    if actual_alignment == expected_alignment:
        score += 0.5
    return score

def impress_slide3_table_at_bottom__db4d3e235d7b0e304c073a25921cbfa2_qw35sft2_ca53a561(result, expected, **options):
    """Check if the table on slide 3 has been moved to the bottom of the slide.

    The table is considered 'at the bottom' when its top edge is at least 60%
    of the slide height from the top (i.e., in the lower 40% of the slide).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    table_top = result.get('table_top')
    slide_height = result.get('slide_height')
    if table_top is None or slide_height is None or slide_height == 0:
        return 0.0
    threshold = expected.get('position_threshold', 0.6)
    min_top = slide_height * threshold
    if table_top >= min_top:
        return 1.0
    return 0.0

def check_slide14_font_sizes__ec3f2f9f447d657b0668843895ae7804_qw35sft2_b758be1f(result, expected, **options):
    """Check that both textboxes on slide 14 have the expected font sizes.
    Expected keys: textbox1_size_pt, textbox2_size_pt.
    Returns 1.0 only if both match; 0.5 if exactly one matches; 0.0 otherwise.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    tb1_ok = result.get('textbox1_size_pt') == expected.get('textbox1_size_pt')
    tb2_ok = result.get('textbox2_size_pt') == expected.get('textbox2_size_pt')
    if tb1_ok:
        score += 0.5
    if tb2_ok:
        score += 0.5
    return score

def check_impress_notes_bg_title__cb759dd1cf89c2aec2802f8543d62400_qw35sft2_1d03f915(result, expected, **options):
    """Check notes text, purple background, and slide title with partial credit.

    Scoring:
      0.33 - notes_text matches expected
      0.34 - background is purple-ish
      0.33 - title_text matches expected_title
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_notes = expected.get('notes_text', '')
    actual_notes = result.get('notes_text', '')
    if actual_notes == expected_notes:
        score += 0.33
    bg_rgb = result.get('bg_rgb')
    if bg_rgb and len(bg_rgb) == 6:
        try:
            r = int(bg_rgb[0:2], 16)
            g = int(bg_rgb[2:4], 16)
            b = int(bg_rgb[4:6], 16)
            if r >= 64 and b >= 64 and (g < 128) and (r - g > 30) and (b - g > 30):
                score += 0.34
        except ValueError:
            pass
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title_text', '')
    if actual_title == expected_title:
        score += 0.33
    return round(score, 4)

def check_pptx_text_color__39fe73dd35955e4e54b2d721092898e1_qw35sft2_d603cba2(result, expected, **options):
    """Check if the extracted text color matches the expected color with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_hex = expected.get('expected_color', '').upper().lstrip('#')
    if not expected_hex or len(expected_hex) != 6:
        return 0.0
    actual_hex = result.get('color_hex', '').upper().lstrip('#')
    if not actual_hex or len(actual_hex) != 6:
        return 0.0
    try:
        tolerance = expected.get('tolerance', 10)
        exp_r = int(expected_hex[0:2], 16)
        exp_g = int(expected_hex[2:4], 16)
        exp_b = int(expected_hex[4:6], 16)
        act_r = result.get('r', 0)
        act_g = result.get('g', 0)
        act_b = result.get('b', 0)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_pptx_text_color__c830024b32b30956e96002af134f6f4a_qw35sft2_a9a956ae(result, expected, **options):
    """Check if the font color of a textbox matches the expected color with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_hex = expected.get('expected_color', '').upper().lstrip('#')
    if not expected_hex:
        return 0.0
    actual_hex = result.get('color_hex', '').upper().lstrip('#')
    if not actual_hex or len(actual_hex) != 6:
        return 0.0
    try:
        exp_r = int(expected_hex[0:2], 16)
        exp_g = int(expected_hex[2:4], 16)
        exp_b = int(expected_hex[4:6], 16)
        act_r = result.get('r', 0)
        act_g = result.get('g', 0)
        act_b = result.get('b', 0)
        tolerance = expected.get('tolerance', 10)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_slide2_bg_subtitle_color__2d907c4b088ee7e1bdf99fbd932517ff_qw35sft2_3aba4d69(result, expected, **options):
    """Check slide 2 background color and subtitle text color. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_bg = (expected.get('bg_color') or '').upper().strip()
    expected_sub = (expected.get('subtitle_color') or '').upper().strip()
    actual_bg = (result.get('bg_color') or '').upper().strip()
    actual_sub = (result.get('subtitle_color') or '').upper().strip()
    if expected_bg and actual_bg == expected_bg:
        score += 0.5
    if expected_sub and actual_sub == expected_sub:
        score += 0.5
    return score

def check_pptx_title_subtitle_font__761afd0460bdb180560f4c7116933937_qw35sft2_8b84cf99(result, expected, **options):
    """Check title text, font name, and subtitle text with partial credit."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('title_text', '')
    actual_text = (result.get('title_text') or '').strip()
    if expected_text and expected_text.lower() == actual_text.lower():
        score += 0.34
    expected_font = expected.get('font_name', '')
    actual_font = (result.get('font_name') or '').strip()
    if expected_font and expected_font.lower() == actual_font.lower():
        score += 0.33
    expected_subtitle = expected.get('subtitle_text', '')
    actual_subtitle = (result.get('subtitle_text') or '').strip()
    if expected_subtitle and expected_subtitle.lower() == actual_subtitle.lower():
        score += 0.33
    return min(score, 1.0)

def check_impress_pptx_props__ed92c6da7770c48f6ebca23c6070cbf5_qw35sft2_55ab6274(result, expected, **options):
    """Check slide 3 Group 6 height == 24cm and slide 6 all textbox fonts == 28pt. Partial credit 0.5+0.5."""
    if isinstance(result, str) or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    expected_height = float(expected.get('slide3_group6_height_cm', 24.0))
    actual_height = result.get('slide3_group6_height_cm')
    if actual_height is not None and abs(actual_height - expected_height) < 0.1:
        score += 0.5
    expected_font = int(expected.get('slide6_font_pt', 28))
    slide6_fonts = result.get('slide6_font_sizes', {})
    if slide6_fonts:
        all_match = all((all((pt == expected_font for pt in pts)) for pts in slide6_fonts.values()))
        if all_match:
            score += 0.5
    return round(score, 4)

def check_pptx_content_and_table__7f43e77e7e3ce83ed182df2ac03d3d3a_qw35sft2_df76f63e(result, expected, **options):
    """Check Slide 2 content contains expected text AND table header cell matches expected value.
    Partial credit: 0.5 per sub-goal.
    expected keys: slide2_content_expected, table_cell_00_expected
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_content = expected.get('slide2_content_expected', '')
    actual_content = result.get('slide2_content', '') if isinstance(result, dict) else ''
    if expected_content and expected_content.lower() in actual_content.lower():
        score += 0.5
    expected_cell = expected.get('table_cell_00_expected', '')
    actual_cell = result.get('table_cell_00', '') if isinstance(result, dict) else ''
    if expected_cell and expected_cell.lower() == actual_cell.lower():
        score += 0.5
    return min(score, 1.0)

def check_picture_size_slide6__f55d52551229ea682c89b1a15474ccff_qw35sft2_54cc8e8d(result, expected, **options):
    """Check picture heights on slides 3,4,6 and width on slide 6. Partial credit 1/4 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    tolerance = 0.5
    score = 0.0
    close_enough = lambda actual, target, tol: actual is not None and abs(actual - target) <= tol
    if close_enough(result.get('slide3_height_cm'), expected.get('slide3_height_cm'), tolerance):
        score += 0.25
    if close_enough(result.get('slide4_height_cm'), expected.get('slide4_height_cm'), tolerance):
        score += 0.25
    if close_enough(result.get('slide6_height_cm'), expected.get('slide6_height_cm'), tolerance):
        score += 0.25
    if close_enough(result.get('slide6_width_cm'), expected.get('slide6_width_cm'), tolerance):
        score += 0.25
    return min(round(score, 4), 1.0)

def check_impress_panel_slide_count__ddf9714ddb153a97bbdfb4350e733817_qw35sft2_992c6a6b(result, expected, **options):
    """
    Partial-credit check:
      0.5 – slide pane is visible (restored)
      0.5 – presentation has at least expected_slide_count slides (default 2)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 2)
    if result.get('panel_visible'):
        score += 0.5
    if result.get('slide_count', 0) >= expected_count:
        score += 0.5
    return score

def check_pptx_font_size__f00d8a2d9828be4b94677a9606cd7f47_qw35sft2_09ce01c6(result, expected, **options):
    """Check that all runs in the target shape have the expected font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_size')
    if expected_size is None:
        return 0.0
    sizes = result.get('font_sizes', [])
    if not sizes:
        return 0.0
    return 1.0 if all((abs(s - expected_size) < 0.5 for s in sizes)) else 0.0

def check_pptx_portrait__2491631e79810e0b815ac32866fd3274_qw35sft2_37505aa0(result, expected, **options):
    """Check that the PPTX slide orientation is portrait (height > width)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('is_portrait', False) else 0.0

def check_pptx_orientation_state__f5c58fc75d059d716d7a9ec9e8e3967d_qw35sft2_4344dc80(result, expected, **options):
    """Partial credit: 0.5 for correct slide count, 0.5 for portrait orientation."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 13)
    if result.get('slide_count') == expected_count:
        score += 0.5
    if result.get('is_portrait') is True:
        score += 0.5
    return min(score, 1.0)

def check_impress_title_at_center__ff47db79abc78d9904bd7e17dbb1e3d3_qw35sft2_e46633fe(result, expected, **options):
    """Check if the slide 2 title has been moved to the vertical center region of the slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    top_cm = result.get('top_cm')
    if top_cm is None:
        return 0.0
    center_min_cm = expected.get('center_min_cm', 5.0)
    center_max_cm = expected.get('center_max_cm', 11.0)
    if center_min_cm <= top_cm <= center_max_cm:
        return 1.0
    return 0.0

def check_pptx_notes_nonempty__d897409e855e2cbbb791ebde03eeef8a_qw35sft2_6c59bb80(result, expected, **options):
    """Check if the slide notes section contains any text (non-empty)."""
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('nonempty') else 0.0

def check_pptx_slide_title_text__c928e38a883d919c295367c551657155_qw35sft2_48305e8b(result, expected, **options):
    """Check if the slide title matches the expected text (case-insensitive)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_title = expected.get('expected_title', '').strip()
    actual_title = result.get('title_text', '').strip()
    if not expected_title:
        return 0.0
    return 1.0 if actual_title.lower() == expected_title.lower() else 0.0

def check_pptx_image_and_transition__4195ceef20cb2389eadfdf0daf18e8e7_qw35sft2_300afbf0(result, expected, **options):
    """Check image size and transition presence with partial credit (0.5 each)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    images = result.get('images', [])
    if len(images) >= expected.get('min_image_count', 1):
        exp_w = float(expected.get('expected_width_cm', 1.0))
        exp_h = float(expected.get('expected_height_cm', 1.0))
        tol = float(expected.get('tolerance_cm', 0.15))
        for img in images:
            if abs(float(img.get('width_cm', 0)) - exp_w) <= tol and abs(float(img.get('height_cm', 0)) - exp_h) <= tol:
                score += 0.5
                break
    if result.get('has_transition', False):
        score += 0.5
    return score

def check_slide_bg_color__b309be1bc27dcbe28bf5749a28cc9b6f_qw35sft2_04bc853b(result, expected, **options):
    """Check that the slide background matches the expected RGB color.

    expected keys:
      - expected_rgb: 6-char hex string, e.g. 'FF6600' for orange
      - tolerance: max Euclidean RGB distance (default 30)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    target_hex = expected.get('expected_rgb', 'FF6600').upper().lstrip('#')
    tolerance = expected.get('tolerance', 30)
    for key in ('rgb', 'master_rgb'):
        actual = result.get(key)
        if actual and len(actual) == 6:
            try:
                r1, g1, b1 = (int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16))
                r2, g2, b2 = (int(target_hex[0:2], 16), int(target_hex[2:4], 16), int(target_hex[4:6], 16))
                dist = ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
            except Exception:
                dist = float('inf')
            if dist <= tolerance:
                return 1.0
    return 0.0

def check_pptx_slide_subtitle_text__4d840f044cdf667c54faf806c6af73e1_qw35sft2_c3b1a416(result, expected, **options):
    """Check if the subtitle text of a slide matches the expected value."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    actual_text = result.get('subtitle_text', '')
    if not actual_text:
        return 0.0
    normalize = lambda s: ' '.join(s.strip().lower().split())
    if normalize(actual_text) == normalize(expected_text):
        return 1.0
    all_texts = result.get('all_texts', [])
    for text in all_texts:
        if normalize(expected_text) in normalize(text) or normalize(text) in normalize(expected_text):
            return 0.8
    return 0.0

def check_pptx_last_two_duplicated__ae132e5d53d956fce4051acdf425fcba_qw35sft2_371de5a2(result, expected, **options):
    """Check that both last two slides were duplicated and appended: total 26 slides."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_count = expected.get('expected_slide_count', 26)
    actual_count = result.get('slide_count', 0)
    if actual_count != expected_count:
        return 0.0
    slide_25_kw = expected.get('slide_25_keyword', 'Now you')
    slide_26_kw = expected.get('slide_26_keyword', 'Works Cited')
    slide_25_text = result.get('slide_25_text', '')
    slide_26_text = result.get('slide_26_text', '')
    score = 0.0
    if slide_25_kw.lower() in slide_25_text.lower():
        score += 0.5
    if slide_26_kw.lower() in slide_26_text.lower():
        score += 0.5
    return score

def check_blank_slide_count__881be08815c16b0f7ba88245e2d69e10_qw35sft2_c7683f4c(result, expected, **options):
    """Check that presentation has expected slide count and all slides are blank (no text frames)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('slide_count')
    if expected_count is not None and result.get('slide_count') == expected_count:
        score += 0.5
    if result.get('slides_with_text_frames', 999) == 0:
        score += 0.5
    return score

def check_impress_audio_and_transition__1b19ca869f4a35a02b8d110c7e17b7e9_qw35sft2_f096de71(result, expected, **options):
    """Check if slide 1 has audio (0.5 credit) and Fade transition (0.5 credit).
    expected keys: expected_transition (str, default 'fade')
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_transition = expected.get('expected_transition', 'fade').lower()
    if result.get('has_audio'):
        score += 0.5
    actual_transition = (result.get('transition_type') or '').lower()
    if actual_transition == expected_transition:
        score += 0.5
    return min(score, 1.0)

def check_pptx_slide_font_italic__5e9f703dc627e70a3b4d09452ec3d3ab_qw35sft2_23ac8c49(result, expected, **options):
    """Check font name (0.5) and italic state (0.5) of last slide title."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('font_name') == expected.get('font_name'):
        score += 0.5
    expected_italic = expected.get('italic')
    result_italic = result.get('italic')
    if expected_italic is True:
        if result_italic is True:
            score += 0.5
    elif result_italic == expected_italic:
        score += 0.5
    return score

def check_pptx_slides35_colors__0fbb4b3515e1009c69bc08919cd2b491_qw35sft2_82a32b53(result, expected, **options):
    """Check text colors: slide 5 should be yellow, slide 3 should be red. Partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    slide5_color = expected.get('slide5_color', 'FFFF00').upper().lstrip('#')
    slide3_color = expected.get('slide3_color', 'FF0000').upper().lstrip('#')
    for slide_key, expected_color in [('slide5', slide5_color), ('slide3', slide3_color)]:
        colors = result.get(slide_key, [])
        if not colors:
            continue
        all_match = all((c == expected_color for c in colors if c is not None))
        none_present = any((c is None for c in colors))
        if all_match and (not none_present):
            score += 0.5
    return min(score, 1.0)

def check_pptx_slide3_subpoint_level__3c50fb39e50d5c4bd155f8f4b7e911d1_qw35sft2_cdee0fb6(result, expected, **options):
    """Check that 'first point of sub topics' is at the expected indentation level."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_level = expected.get('expected_level', 0)
    actual_level = result.get('level', -1)
    return 1.0 if actual_level == expected_level else 0.0

def check_save_and_slide_count__e4d3506bfaa42b2028f1a93213c30bc6_qw35sft2_73a87272(result, expected, **options):
    """Check that pre.pptx was saved on Desktop (0.5) and has expected slide count (0.5)."""
    if isinstance(result, str) or result is None:
        return 0.0
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
    expected_count = expected.get('expected_slide_count', 18)
    actual_count = result.get('slide_count', 0)
    if actual_count == expected_count:
        score += 0.5
    return min(score, 1.0)

def check_impress_bullet_underline__396e098239a8be5c53a96d982334cb7b_qw35sft2_6b235811(result, expected, **options):
    """Check bullet + underline: 0.5 for bullet, 0.5 for underline."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_bullet') is True:
        score += 0.5
    if result.get('has_underline') is True:
        score += 0.5
    return score

def check_impress_options_combo__43da862b24de1bb5002a875bb93d754e_qw35sft2_6c066c1d(result, expected, **options):
    """
    Check both: presenter console disabled (0.5) + autosave interval matches (0.5).
    expected keys: autosave_minutes (int)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('presenter_disabled', False):
        score += 0.5
    expected_minutes = expected.get('autosave_minutes')
    actual_minutes = result.get('autosave_minutes')
    if expected_minutes is not None and actual_minutes == expected_minutes:
        score += 0.5
    return score

def check_pptx_slide_orientation__1830f729a2ae0637260a8bca58b0c8cd_qw35sft2_c6ecb7c1(result, expected, **options):
    """Check if the slide orientation matches the expected value ('portrait' or 'landscape')."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_orientation = expected.get('expected_orientation', 'portrait').lower()
    actual_orientation = result.get('orientation', '').lower()
    return 1.0 if actual_orientation == expected_orientation else 0.0

def check_impress_img_top__edae955a62ba6a9d5bff7ce774e3ee10_qw35sft2_5c366587(result, expected, **options):
    """
    Check that the picture on slide 2 has been moved to the slide top region.
    'top_threshold_emu' in expected defines the maximum top value (in EMU) to be considered 'at top'.
    Default: 1638000 EMU (~4.55 cm), well above the original position ~6.94 cm (2499480 EMU).
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    image_top = result.get('image_top_emu')
    if image_top is None:
        return 0.0
    threshold = expected.get('top_threshold_emu', 1638000)
    return 1.0 if image_top <= threshold else 0.0

def check_impress_title_color_underline__84ce98fac5a701a1ff43f1337a4f838c_qw35sft2_b63365eb(result, expected, **options):
    """
    Check if the title has the expected color and underline (variation 4, slide 3).
    expected keys: expected_color_hex (e.g. '000000'), expected_underline (True)
    Partial credit: 0.5 per criterion.
    Returns: 0.0, 0.5, or 1.0
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_color = str(expected.get('expected_color_hex', '000000')).upper().lstrip('#')
    all_colors = result.get('all_colors', [result.get('color_hex')])
    if all_colors and all((c is not None and c.upper().lstrip('#') == expected_color for c in all_colors)):
        score += 0.5
    expected_underline = expected.get('expected_underline', True)
    if result.get('underline') == expected_underline:
        score += 0.5
    return score

def check_pptx_slide2_title_bold__ba024eb57784584ce2b08e31db5871ac_qw35sft2_6b563f6f(result, expected, **options):
    """
    Partial credit:
    - 0.34: title text equals expected
    - 0.33: alignment is RIGHT
    - 0.33: title is bold
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('title_text') == expected.get('title_text', 'Note'):
        score += 0.34
    if result.get('title_alignment') == expected.get('title_alignment', 'RIGHT'):
        score += 0.33
    if result.get('title_bold') == expected.get('title_bold', True):
        score += 0.33
    return min(score, 1.0)

def check_impress_slide3_title_bold__7dcab435f4cba9d9e23699eed9a4d408_qw35sft2_1080995a(result, expected, **options):
    """Check if the title in slide 3 is bold."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    title_bold = result.get('title_bold')
    if title_bold is True:
        return 1.0
    return 0.0

def check_pptx_text_alignments__14ff9da6ca84bfee11448e7c7f8efafc_qw35sft2_86d84f40(result, expected, **options):
    """Check text alignment for slides 3, 4, 5. Partial credit: 1/3 per slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    per_slide = 1.0 / 3
    if result.get('slide3_align') == expected.get('slide3_align'):
        score += per_slide
    if result.get('slide4_align') == expected.get('slide4_align'):
        score += per_slide
    if result.get('slide5_align') == expected.get('slide5_align'):
        score += per_slide
    return min(round(score, 4), 1.0)

def check_impress_notes_bg_transition__5cb083f40f95f75316141f060ee8a1cc_qw35sft2_5f8249e7(result, expected, **options):
    """Check notes text, purple background, and slide transition with partial credit.

    Scoring:
      0.33 - notes_text matches expected
      0.34 - background is purple-ish
      0.33 - slide has a transition set
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_notes = expected.get('notes_text', '')
    actual_notes = result.get('notes_text', '')
    if actual_notes == expected_notes:
        score += 0.33
    bg_rgb = result.get('bg_rgb')
    if bg_rgb and len(bg_rgb) == 6:
        try:
            r = int(bg_rgb[0:2], 16)
            g = int(bg_rgb[2:4], 16)
            b = int(bg_rgb[4:6], 16)
            if r >= 64 and b >= 64 and (g < 128) and (r - g > 30) and (b - g > 30):
                score += 0.34
        except ValueError:
            pass
    if result.get('has_transition') is True:
        score += 0.33
    return round(score, 4)

def check_slide14_font_bold__a28710b7c875e09d23efa44313d2c747_qw35sft2_644b3032(result, expected, **options):
    """Check slide 14: textbox1 size, textbox1 bold, textbox2 size.
    Partial credit: 0.4 for tb1 size, 0.2 for tb1 bold, 0.4 for tb2 size.
    Expected keys: textbox1_size_pt, textbox1_bold, textbox2_size_pt.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('textbox1_size_pt') == expected.get('textbox1_size_pt'):
        score += 0.4
    expected_bold = expected.get('textbox1_bold', True)
    actual_bold = result.get('textbox1_bold')
    if expected_bold:
        if actual_bold is True:
            score += 0.2
    elif actual_bold is not True:
        score += 0.2
    if result.get('textbox2_size_pt') == expected.get('textbox2_size_pt'):
        score += 0.4
    return min(score, 1.0)

def impress_slide3_bottom_and_title__80d9e8d7efd024cbc51991913269b5fb_qw35sft2_fba30bc1(result, expected, **options):
    """Check that the table is at the bottom of slide 3 AND the slide title was renamed.

    Partial credit:
      - 0.5 for table position at bottom (top >= 60% of slide height)
      - 0.5 for slide 3 title matching expected_title
    """
    if isinstance(result, dict) and result.get('error') and (result.get('table_top') is None):
        return 0.0
    score = 0.0
    table_top = result.get('table_top')
    slide_height = result.get('slide_height')
    if table_top is not None and slide_height and (slide_height > 0):
        threshold = expected.get('position_threshold', 0.6)
        if table_top >= slide_height * threshold:
            score += 0.5
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title') or ''
    if expected_title and actual_title.strip() == expected_title.strip():
        score += 0.5
    return min(score, 1.0)

def check_impress_slide_title_full__9620e3371e8842ebb86c9c7ef6cb4f00_qw35sft2_fa068e33(result, expected, **options):
    """Check slide title text (0.34), font size (0.33), and font color (0.33) with partial credit."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('expected_title', '')
    actual_text = result.get('title_text', '')
    if expected_text and actual_text.strip() == expected_text.strip():
        score += 0.34
    expected_size = expected.get('expected_font_size_pt')
    actual_size = result.get('font_size_pt')
    size_tolerance = expected.get('size_tolerance', 1.0)
    if expected_size is not None and actual_size is not None:
        if abs(actual_size - expected_size) <= size_tolerance:
            score += 0.33
    expected_hex = expected.get('expected_color', '').upper().lstrip('#')
    actual_hex = result.get('color_hex', '')
    if actual_hex:
        actual_hex = actual_hex.upper().lstrip('#')
    if expected_hex and actual_hex and (len(expected_hex) == 6) and (len(actual_hex) == 6):
        try:
            exp_r = int(expected_hex[0:2], 16)
            exp_g = int(expected_hex[2:4], 16)
            exp_b = int(expected_hex[4:6], 16)
            act_r = result.get('r', 0)
            act_g = result.get('g', 0)
            act_b = result.get('b', 0)
            tolerance = expected.get('tolerance', 10)
            if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
                score += 0.33
        except Exception:
            pass
    return min(score, 1.0)

def check_pptx_title_font_transition__7950498f55dec024fcd49c95968dd860_qw35sft2_2c7c0d01(result, expected, **options):
    """Check title text, font name, and slide transition type with partial credit."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_text = expected.get('title_text', '')
    actual_text = (result.get('title_text') or '').strip()
    if expected_text and expected_text.lower() == actual_text.lower():
        score += 0.34
    expected_font = expected.get('font_name', '')
    actual_font = (result.get('font_name') or '').strip()
    if expected_font and expected_font.lower() == actual_font.lower():
        score += 0.33
    expected_transition = expected.get('transition_type', '').lower()
    actual_transition = (result.get('transition_type') or '').lower()
    if expected_transition and actual_transition and (expected_transition == actual_transition):
        score += 0.33
    return min(score, 1.0)

def check_pptx_text_color__f2292754a9e180ffad3521e45a75c905_qw35sft2_dc1d29ca(result, expected, **options):
    """Check if the extracted text color matches the expected color with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_hex = expected.get('expected_color', '').upper().lstrip('#')
    if not expected_hex or len(expected_hex) != 6:
        return 0.0
    actual_hex = result.get('color_hex', '').upper().lstrip('#')
    if not actual_hex or len(actual_hex) != 6:
        return 0.0
    try:
        tolerance = expected.get('tolerance', 10)
        exp_r = int(expected_hex[0:2], 16)
        exp_g = int(expected_hex[2:4], 16)
        exp_b = int(expected_hex[4:6], 16)
        act_r = result.get('r', 0)
        act_g = result.get('g', 0)
        act_b = result.get('b', 0)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_slide2_bg_title_bold__3ee304ba526b814e3d4d6c2be946eae0_qw35sft2_0f122ed4(result, expected, **options):
    """Check slide 2 background color and title bold status. Partial credit 0.5 each."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_bg = (expected.get('bg_color') or '').upper().strip()
    expected_bold = expected.get('title_bold', True)
    actual_bg = (result.get('bg_color') or '').upper().strip()
    actual_bold = result.get('title_bold')
    if expected_bg and actual_bg == expected_bg:
        score += 0.5
    if actual_bold is True and expected_bold is True:
        score += 0.5
    return score

def check_pptx_text_color__c8d6e08f7da191ce71fc333f5a053d63_qw35sft2_1bd747b0(result, expected, **options):
    """Check if the font color of a textbox matches the expected color with tolerance."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_hex = expected.get('expected_color', '').upper().lstrip('#')
    if not expected_hex:
        return 0.0
    actual_hex = result.get('color_hex', '').upper().lstrip('#')
    if not actual_hex or len(actual_hex) != 6:
        return 0.0
    try:
        exp_r = int(expected_hex[0:2], 16)
        exp_g = int(expected_hex[2:4], 16)
        exp_b = int(expected_hex[4:6], 16)
        act_r = result.get('r', 0)
        act_g = result.get('g', 0)
        act_b = result.get('b', 0)
        tolerance = expected.get('tolerance', 10)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_pptx_content_bold__87b69c70091c302ea2d6ddb5f5d9c002_qw35sft2_5ef7630d(result, expected, **options):
    """Check Slide 2 content contains expected text AND the text uses bold formatting.
    Partial credit: 0.5 per sub-goal.
    expected keys: slide2_content_expected, slide2_content_bold_expected
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_content = expected.get('slide2_content_expected', '')
    actual_content = result.get('slide2_content', '') if isinstance(result, dict) else ''
    if expected_content and expected_content.lower() in actual_content.lower():
        score += 0.5
    expected_bold = expected.get('slide2_content_bold_expected', True)
    actual_bold = result.get('slide2_content_bold', False) if isinstance(result, dict) else False
    if actual_bold == expected_bold:
        score += 0.5
    return min(score, 1.0)

def check_impress_slide_pane__0dbd76d4f35e1d6b60aa20d8163a50c2_qw35sft2_d7e6c0f9(result, expected, **options):
    """Check that the left slide pane ('Slides' panel) is visible in LibreOffice Impress."""
    if not isinstance(result, dict):
        return 0.0
    tree = result.get('tree', '')
    if 'Slides' in tree:
        return 1.0
    return 0.0

def check_pptx_font_size__8cd3b2c6f2a2d80e7aa650d3b0962695_qw35sft2_700a1334(result, expected, **options):
    """Check that all runs in the target shape have the expected font size."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_size = expected.get('expected_size')
    if expected_size is None:
        return 0.0
    sizes = result.get('font_sizes', [])
    if not sizes:
        return 0.0
    return 1.0 if all((abs(s - expected_size) < 0.5 for s in sizes)) else 0.0

def check_impress_pptx_props__e526a93f073489d810264dc88333d90d_qw35sft2_f459912c(result, expected, **options):
    """Check slide 3 Group 6 height == 18cm and slide 6 all textbox fonts == 36pt. Partial credit 0.5+0.5."""
    if isinstance(result, str) or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    expected_height = float(expected.get('slide3_group6_height_cm', 18.0))
    actual_height = result.get('slide3_group6_height_cm')
    if actual_height is not None and abs(actual_height - expected_height) < 0.1:
        score += 0.5
    expected_font = int(expected.get('slide6_font_pt', 36))
    slide6_fonts = result.get('slide6_font_sizes', {})
    if slide6_fonts:
        all_match = all((all((pt == expected_font for pt in pts)) for pts in slide6_fonts.values()))
        if all_match:
            score += 0.5
    return round(score, 4)

def check_pptx_has_transition__24d25146d014952eed8e287afb1ec612_qw35sft2_3e97889e(result, expected, **options):
    """Check that a transition effect has been applied to the slide."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    return 1.0 if result.get('has_transition', False) else 0.0

def check_pptx_transitions__b146bd229b2aa1513d722fe51dc492dc_qw35sft2_1a2a7e26(result, expected, **options):
    """Check that specific slides have the expected transition types. Partial credit per slide."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    score = 0.0
    per_check = 1.0 / len(checks)
    for check in checks:
        idx = check['slide_idx']
        expected_type = check['transition_type']
        actual = result.get('slide_{}'.format(idx))
        if actual == expected_type:
            score += per_check
    return min(round(score, 4), 1.0)

def check_pptx_table_row0__a30c5776ecd0bc922b282656ddbb8d0e_qw35sft2_978d7b1e(result, expected, **options):
    """Check if the first row of the table matches expected values. All 4 cells must match."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_row = expected.get('expected_table_row0', [])
    actual_row = result.get('table_row0', [])
    if not expected_row or not actual_row:
        return 0.0
    if len(actual_row) != len(expected_row):
        return 0.0
    matches = sum((1 for a, e in zip(actual_row, expected_row) if a == e))
    return matches / len(expected_row)

def check_pptx_slide_title_text__a33045396be2560e0d47268bd7a68b1f_qw35sft2_af87e143(result, expected, **options):
    """Check if the slide title matches the expected text (case-insensitive)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_title = expected.get('expected_title', '').strip()
    actual_title = result.get('title_text', '').strip()
    if not expected_title:
        return 0.0
    return 1.0 if actual_title.lower() == expected_title.lower() else 0.0

def check_slide_bg_color__3e8a909d4fdee8dfddabe4251505a445_qw35sft2_3ed277f0(result, expected, **options):
    """Check that the slide background matches the expected RGB color.

    expected keys:
      - expected_rgb: 6-char hex string, e.g. 'FF0000'
      - tolerance: max Euclidean RGB distance (default 30)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    target_hex = expected.get('expected_rgb', 'FF0000').upper().lstrip('#')
    tolerance = expected.get('tolerance', 30)
    for key in ('rgb', 'master_rgb'):
        actual = result.get(key)
        if actual and len(actual) == 6:
            try:
                r1, g1, b1 = (int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16))
                r2, g2, b2 = (int(target_hex[0:2], 16), int(target_hex[2:4], 16), int(target_hex[4:6], 16))
                dist = ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5
            except Exception:
                dist = float('inf')
            if dist <= tolerance:
                return 1.0
    return 0.0

def check_blank_slide_count__b298657a7c2d705873344441c4e6373c_qw35sft2_6a7840bc(result, expected, **options):
    """Check that presentation has expected slide count and all slides are blank (no text frames)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('slide_count')
    if expected_count is not None and result.get('slide_count') == expected_count:
        score += 0.5
    if result.get('slides_with_text_frames', 999) == 0:
        score += 0.5
    return score

def check_pptx_image_and_title__68b16c8fa261fe6f7f7bf99cd53e207c_qw35sft2_a9adceaf(result, expected, **options):
    """Check image size and title text with partial credit (0.5 each)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    images = result.get('images', [])
    if len(images) >= expected.get('min_image_count', 1):
        exp_w = float(expected.get('expected_width_cm', 1.5))
        exp_h = float(expected.get('expected_height_cm', 1.5))
        tol = float(expected.get('tolerance_cm', 0.15))
        for img in images:
            if abs(float(img.get('width_cm', 0)) - exp_w) <= tol and abs(float(img.get('height_cm', 0)) - exp_h) <= tol:
                score += 0.5
                break
    expected_title = expected.get('expected_title', '')
    actual_title = result.get('title_text') or ''
    if expected_title and expected_title.strip().lower() == actual_title.strip().lower():
        score += 0.5
    return score

def check_impress_title_bottom_and_font__8ff66c4d1cd7e4775354b5fdf25e12eb_qw35sft2_d3fefb96(result, expected, **options):
    """Partial credit: 0.5 for title at bottom, 0.5 for font size == 32pt."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    top_cm = result.get('top_cm')
    min_top_cm = expected.get('min_top_cm', 9.525)
    if top_cm is not None and top_cm >= min_top_cm:
        score += 0.5
    font_size_pt = result.get('font_size_pt')
    expected_font_pt = expected.get('expected_font_pt', 32.0)
    if font_size_pt is not None and abs(font_size_pt - expected_font_pt) < 0.5:
        score += 0.5
    return min(score, 1.0)

def check_pptx_notes_title__e6eecf5158b3e658cb4aa48b9b009898_qw35sft2_e6001705(result, expected, **options):
    """Check if the slide notes contain the expected title text (case-insensitive substring match)."""
    if result.get('error'):
        return 0.0
    notes = result.get('notes', '')
    if not notes:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 1.0 if notes else 0.0
    if expected_text.lower() in notes.lower():
        return 1.0
    return 0.0

def check_two_slide_tables__422453ea9133a41f5d3b73c5d61c25f4_qw35sft2_1314e343(result, expected, **options):
    """Check tables on two slides with partial credit.
    0.5 for Features slide (5 rows x 2 cols), 0.5 for Product Overview slide (2 rows x 3 cols).
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    features_rows = expected.get('features_rows', 5)
    features_cols = expected.get('features_cols', 2)
    slide3_tables = result.get('slide3_tables', [])
    for tbl in slide3_tables:
        if tbl.get('rows') == features_rows and tbl.get('cols') == features_cols:
            score += 0.5
            break
    overview_rows = expected.get('overview_rows', 2)
    overview_cols = expected.get('overview_cols', 3)
    slide2_tables = result.get('slide2_tables', [])
    for tbl in slide2_tables:
        if tbl.get('rows') == overview_rows and tbl.get('cols') == overview_cols:
            score += 0.5
            break
    return min(score, 1.0)

def check_pptx_slide5_text_colors__7491f6c65f733a38f893bdf8cad068b0_qw35sft2_659f48ec(result, expected, **options):
    """Check that all text runs in slide 5 have the expected font color."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    colors = result.get('colors', [])
    if not colors:
        return 0.0
    expected_color = expected.get('expected_color', 'FF0000').upper().lstrip('#')
    all_match = all((c == expected_color for c in colors if c is not None))
    none_present = any((c is None for c in colors))
    if all_match and (not none_present):
        return 1.0
    return 0.0

def check_impress_has_audio__5688a332db768af3543ab15e11555c76_qw35sft2_138ab214(result, expected, **options):
    """Check if the presentation slide has an embedded audio file.
    Returns 1.0 if audio is found, 0.0 otherwise.
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    return 1.0 if result.get('has_audio') else 0.0

def check_pptx_multi_slide_bg_colors__01d7b3b210af327bea300a303c9bddd6_qw35sft2_9c2a1925(result, expected, **options):
    """Check if multiple slides have yellow background colors (partial credit per slide)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    slide_colors = result.get('slide_colors', {})
    if not slide_colors:
        return 0.0
    expected_color = expected.get('expected_color', 'FFFF00').upper().lstrip('#')
    tolerance = expected.get('tolerance', 30)
    slide_indices = expected.get('slide_indices', ['2', '3'])
    try:
        exp_r = int(expected_color[0:2], 16)
        exp_g = int(expected_color[2:4], 16)
        exp_b = int(expected_color[4:6], 16)
    except Exception:
        return 0.0
    score = 0.0
    credit_per_slide = 1.0 / len(slide_indices)
    for idx in slide_indices:
        idx_str = str(idx)
        color_info = slide_colors.get(idx_str, {})
        if color_info.get('error') or color_info.get('color_hex') is None:
            continue
        act_r = color_info.get('r', 0)
        act_g = color_info.get('g', 0)
        act_b = color_info.get('b', 0)
        if abs(act_r - exp_r) <= tolerance and abs(act_g - exp_g) <= tolerance and (abs(act_b - exp_b) <= tolerance):
            score += credit_per_slide
    return min(score, 1.0)

def check_impress_video_docs_vlc__bae283564a19c9e043588bc4818877d3_qw35sft2_c7d97e82(result, expected, **options):
    """Score: 0.5 for video file in ~/Documents/, 0.5 for VLC running."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('video_in_documents'):
        score += 0.5
    if result.get('vlc_running'):
        score += 0.5
    return min(score, 1.0)

def check_impress_compose4__c60351f10d67b0cf97ff8e001ca27d87_qw35sft2_9a3a7691(result, expected, **options):
    """
    Partial credit:
      0.5 - background.png exists on Desktop
      0.5 - PPTX slide 2 image has brightness (a:lum) attribute stored,
            with value >= expected_brightness_pct (e.g. 25%)
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    min_brightness_pct = expected.get('min_brightness_pct', 50.0)
    if result.get('file_exists'):
        score += 0.5
    if result.get('pptx_has_brightness'):
        bv = result.get('brightness_value')
        if bv is not None and bv >= min_brightness_pct:
            score += 0.5
        elif bv is None:
            score += 0.25
    return min(score, 1.0)

def check_pptx_slide2_bg_and_count__24ae1553823b67948a98514a9ed74ad3_qw35sft2_53ba711c(result, expected, **options):
    """Check that slide 2 has a picture background and the presentation still has 20 slides."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_slide_count', 20)
    if result.get('has_background_image'):
        score += 0.7
    if result.get('slide_count') == expected_count:
        score += 0.3
    return min(score, 1.0)

def check_pptx_notes_and_docx__1f2059f6b413a703d90fb8946ec0e698_qw35sft2_3c27e232(result, expected, **options):
    """
    Partial credit scoring:
    - 0.5: slide 7 in PPTX has the expected note
    - 0.5: docx contains the slide 7 note in extracted content
    """
    if not result or (result.get('errors') and len(result['errors']) >= 2):
        return 0.0
    score = 0.0
    expected_note = expected.get('slide7_note', 'No content slide.')
    pptx_note = result.get('pptx_slide7_note')
    if pptx_note is not None and expected_note.strip().lower() in pptx_note.strip().lower():
        score += 0.5
    docx_lines = result.get('docx_lines', [])
    for line in docx_lines:
        if expected_note.strip().lower() in line.strip().lower():
            score += 0.5
            break
    return round(min(score, 1.0), 2)

def check_pptx_slide2_bg__d5c6210d63b8dbe359320c4bfae3310a_qw35sft2_7e1c3639(result, expected, **options):
    """Check that slide 2 has a picture background set."""
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    if result.get('has_background_image'):
        return 1.0
    if result.get('has_any_bg_fill'):
        return 0.5
    return 0.0

def check_impress_compose0__9880fcd0ac1db0e1d57c564194f1780a_qw35sft2_0c53865a(result, expected, **options):
    """
    Partial credit:
      0.5 - background.png exists on Desktop AND image brightness ~50% verified in PPTX XML
      0.5 - slide 2 title matches expected_title
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_title = expected.get('expected_title', '')
    if result.get('file_exists') and result.get('brightness_ok'):
        score += 0.5
    actual_title = result.get('slide2_title') or ''
    if expected_title and actual_title.strip().lower() == expected_title.strip().lower():
        score += 0.5
    return min(score, 1.0)

def check_vlc_dark_slider_oneinstance__d71f004306b232a878090dc93cee8c3a_qw35sft2_07da77f7(result, expected, **options):
    """
    Check two VLC settings:
    1. qt-slider-colours is a blackish color (all RGB channels < 100) - 0.5 credit
    2. one-instance-when-started-from-file equals expected_one_instance (default 1; disable = 0) - 0.5 credit
    """
    score = 0.0
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_9a65d1.error(f"Cannot read VLC config file at '{result}': {e}")
        return 0.0
    config = {}
    for line in config_file.split('\n'):
        stripped = line.strip()
        if not stripped or stripped.startswith('#') or '=' not in stripped:
            continue
        key, _, val = stripped.partition('=')
        config[key.strip()] = val.strip()
    qt_slider = config.get('qt-slider-colours', '153;210;153;20;210;20;255;199;15;245;39;29')
    try:
        values = [int(x) for x in qt_slider.split(';')]
        if len(values) >= 3:
            colors = list(zip(values[0::3], values[1::3], values[2::3]))
            if colors and all((all((c < 100 for c in color)) for color in colors)):
                score += 0.5
                logger_qw35sft2_9a65d1.info(f'Volume slider blackish check passed: {qt_slider}')
            else:
                logger_qw35sft2_9a65d1.info(f'Volume slider not blackish: {colors}')
    except Exception as e:
        logger_qw35sft2_9a65d1.error(f"Error parsing qt-slider-colours '{qt_slider}': {e}")
    expected_one = str(expected.get('expected_one_instance', '0'))
    actual_one = config.get('one-instance-when-started-from-file', '1')
    if actual_one == expected_one:
        score += 0.5
        logger_qw35sft2_9a65d1.info(f'one-instance check passed: {actual_one} == {expected_one}')
    else:
        logger_qw35sft2_9a65d1.info(f'one-instance check failed: {actual_one} != {expected_one}')
    return min(score, 1.0)

def check_vlc_dark_slider_minimal__d20b1f499db3b1fe73151927397d72b8_qw35sft2_184f1657(result, expected, **options):
    """
    Check two VLC Qt settings:
    1. qt-slider-colours is a blackish color (all RGB channels < 100) - 0.5 credit
    2. qt-minimal-view equals expected_minimal_view (default 0; enable = 1) - 0.5 credit
    """
    score = 0.0
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_8ff4c2.error(f"Cannot read VLC config file at '{result}': {e}")
        return 0.0
    config = {}
    for line in config_file.split('\n'):
        stripped = line.strip()
        if not stripped or stripped.startswith('#') or '=' not in stripped:
            continue
        key, _, val = stripped.partition('=')
        config[key.strip()] = val.strip()
    qt_slider = config.get('qt-slider-colours', '153;210;153;20;210;20;255;199;15;245;39;29')
    try:
        values = [int(x) for x in qt_slider.split(';')]
        if len(values) >= 3:
            colors = list(zip(values[0::3], values[1::3], values[2::3]))
            if colors and all((all((c < 100 for c in color)) for color in colors)):
                score += 0.5
                logger_qw35sft2_8ff4c2.info(f'Volume slider blackish check passed: {qt_slider}')
            else:
                logger_qw35sft2_8ff4c2.info(f'Volume slider not blackish: {colors}')
    except Exception as e:
        logger_qw35sft2_8ff4c2.error(f"Error parsing qt-slider-colours '{qt_slider}': {e}")
    expected_minimal = str(expected.get('expected_minimal_view', '1'))
    actual_minimal = config.get('qt-minimal-view', '0')
    if actual_minimal == expected_minimal:
        score += 0.5
        logger_qw35sft2_8ff4c2.info(f'minimal-view check passed: {actual_minimal} == {expected_minimal}')
    else:
        logger_qw35sft2_8ff4c2.info(f'minimal-view check failed: {actual_minimal} != {expected_minimal}')
    return min(score, 1.0)

def check_vlc_dark_slider_maxvol__125763cfb25260df78f204fa214a8b29_qw35sft2_dd32e44e(result, expected, **options):
    """
    Check two VLC Qt settings:
    1. qt-slider-colours is a blackish color (all RGB channels < 100) - 0.5 credit
    2. qt-max-volume equals expected_max_volume - 0.5 credit
    """
    score = 0.0
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_2c2c74.error(f"Cannot read VLC config file at '{result}': {e}")
        return 0.0
    config = {}
    for line in config_file.split('\n'):
        stripped = line.strip()
        if not stripped or stripped.startswith('#') or '=' not in stripped:
            continue
        key, _, val = stripped.partition('=')
        config[key.strip()] = val.strip()
    qt_slider = config.get('qt-slider-colours', '153;210;153;20;210;20;255;199;15;245;39;29')
    try:
        values = [int(x) for x in qt_slider.split(';')]
        if len(values) >= 3:
            colors = list(zip(values[0::3], values[1::3], values[2::3]))
            if colors and all((all((c < 100 for c in color)) for color in colors)):
                score += 0.5
                logger_qw35sft2_2c2c74.info(f'Volume slider blackish check passed: {qt_slider}')
            else:
                logger_qw35sft2_2c2c74.info(f'Volume slider not blackish: {colors}')
    except Exception as e:
        logger_qw35sft2_2c2c74.error(f"Error parsing qt-slider-colours '{qt_slider}': {e}")
    expected_max = str(expected.get('expected_max_volume', '200'))
    actual_max = config.get('qt-max-volume', '125')
    if actual_max == expected_max:
        score += 0.5
        logger_qw35sft2_2c2c74.info(f'Max volume check passed: {actual_max} == {expected_max}')
    else:
        logger_qw35sft2_2c2c74.info(f'Max volume check failed: {actual_max} != {expected_max}')
    return min(score, 1.0)

def check_vlc_dark_slider_bgcone__914b13738d8dcf766e175adf76980e8b_qw35sft2_0920035b(result, expected, **options):
    """
    Check two VLC Qt settings:
    1. qt-slider-colours is a blackish color (all RGB channels < 100) - 0.5 credit
    2. qt-bgcone equals expected_bgcone (default is 1; disable = 0) - 0.5 credit
    """
    score = 0.0
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_256b9d.error(f"Cannot read VLC config file at '{result}': {e}")
        return 0.0
    config = {}
    for line in config_file.split('\n'):
        stripped = line.strip()
        if not stripped or stripped.startswith('#') or '=' not in stripped:
            continue
        key, _, val = stripped.partition('=')
        config[key.strip()] = val.strip()
    qt_slider = config.get('qt-slider-colours', '153;210;153;20;210;20;255;199;15;245;39;29')
    try:
        values = [int(x) for x in qt_slider.split(';')]
        if len(values) >= 3:
            colors = list(zip(values[0::3], values[1::3], values[2::3]))
            if colors and all((all((c < 100 for c in color)) for color in colors)):
                score += 0.5
                logger_qw35sft2_256b9d.info(f'Volume slider blackish check passed: {qt_slider}')
            else:
                logger_qw35sft2_256b9d.info(f'Volume slider not blackish: {colors}')
    except Exception as e:
        logger_qw35sft2_256b9d.error(f"Error parsing qt-slider-colours '{qt_slider}': {e}")
    expected_bgcone = str(expected.get('expected_bgcone', '0'))
    actual_bgcone = config.get('qt-bgcone', '1')
    if actual_bgcone == expected_bgcone:
        score += 0.5
        logger_qw35sft2_256b9d.info(f'bgcone check passed: {actual_bgcone} == {expected_bgcone}')
    else:
        logger_qw35sft2_256b9d.info(f'bgcone check failed: {actual_bgcone} != {expected_bgcone}')
    return min(score, 1.0)

def check_vlc_dark_slider_globalkey__c70b3a78c573774afb864d5784b3c8df_qw35sft2_56f9e06c(result, expected, **options):
    """
    Check two VLC settings:
    1. qt-slider-colours is a blackish color (all RGB channels < 100) - 0.5 credit
    2. global-key-play-pause is set to a non-empty value (expected_global_key = "1") - 0.5 credit

    For the global hotkey check: "0" means not set (empty value), "1" means set.
    """
    score = 0.0
    try:
        with open(result, 'rb') as f:
            config_file = f.read().decode('utf-8')
    except Exception as e:
        logger_qw35sft2_589392.error(f"Cannot read VLC config file at '{result}': {e}")
        return 0.0
    config = {}
    for line in config_file.split('\n'):
        stripped = line.strip()
        if not stripped or stripped.startswith('#') or '=' not in stripped:
            continue
        key, _, val = stripped.partition('=')
        config[key.strip()] = val.strip()
    qt_slider = config.get('qt-slider-colours', '153;210;153;20;210;20;255;199;15;245;39;29')
    try:
        values = [int(x) for x in qt_slider.split(';')]
        if len(values) >= 3:
            colors = list(zip(values[0::3], values[1::3], values[2::3]))
            if colors and all((all((c < 100 for c in color)) for color in colors)):
                score += 0.5
                logger_qw35sft2_589392.info(f'Volume slider blackish check passed: {qt_slider}')
            else:
                logger_qw35sft2_589392.info(f'Volume slider not blackish: {colors}')
    except Exception as e:
        logger_qw35sft2_589392.error(f"Error parsing qt-slider-colours '{qt_slider}': {e}")
    expected_global_key = str(expected.get('expected_global_key', '1'))
    raw_val = config.get('global-key-play-pause', '')
    actual_global_key = '0' if raw_val == '' else '1'
    if actual_global_key == expected_global_key:
        score += 0.5
        logger_qw35sft2_589392.info(f"global-key-play-pause check passed: raw='{raw_val}', interpreted={actual_global_key}")
    else:
        logger_qw35sft2_589392.info(f"global-key-play-pause check failed: raw='{raw_val}', interpreted={actual_global_key} != {expected_global_key}")
    return min(score, 1.0)
