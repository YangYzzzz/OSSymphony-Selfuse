"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_cell_value__3a7517792dc1ec4e3be4f857bdb0950f', 'check_xlsx_reversed_names__d8dd1c41700e60d858e6fd9c2f1f5451', 'check_cell_approx__3440d9b31b1654ffe8da4da80be4e5fc', 'check_cell_numeric_value__db45e57fbcc1253e5b0f6029b5341c3e', 'check_ods_to_xlsx__d8b4b0a7ce586e5e25f592627cc186b5', 'check_cell_numeric__f3f19b96be4e433759cb0f99e36be750', 'check_cell_b3_value__0698f7f5da73d25e009ccfe66a594e11', 'check_cell_value__4a0f3185fd3c30bf019ccadcf5573a11', 'check_multi_cells__677226b05224be2ecbd2de2e693f9c8e', 'check_cell_sum_value__dc007ee2aea5c445c94303ee0d01fcf4', 'check_filtered_sheet__0d0f959234d0c494b6d7c17cb5158836', 'check_xlsx_new_row__246ef0821ed926b9e34f1369cfdfe795', 'check_range_cells__980775095293f828bec4466da30db27b', 'check_multi_cell_values__6d0d805c4a52313050aa161e6524aaf8', 'check_cell_values_v0__c8482c50afd0db1eeb95cd8ec431bb1e', 'check_xlsx_cell_value__fe91ae8d0e7c34f2eef1f4db512bed2a', 'check_sheet_names__9889ac58b2fb5d2c955b2821b86f86b7', 'check_multi_cell_values__8a05426d089eda14388178863fc78e0b', 'check_sheet_name__e743a4ede9bcdbd814ac13f073acba4a', 'check_sheet_renamed__64efe2d9a5895dad59edc0eb16307294', 'check_multi_cell__1393f3a73542b23a8e048e3ef0292819', 'check_cell_value__3ad3a7a9e15d5bbcab97b429c7a2d70d', 'check_xlsx_cell_value__5729380214ac80099631e898f8c00bcc', 'check_cell_numeric_value__c13e5303d821a816dfe6dbe089834e52', 'check_cell_value_numeric__ad5647b3d1b47a9e1996a6e69ca9562d', 'check_cell_numeric_value__67a1f86b27a24adb60e8a183c35b0852', 'check_cells_replaced__1cd51ca9c28aca9840aa590661974218', 'check_xlsx_filenames__aad13293a214d4a4614be1039d459c78', 'check_cell_string_value__69f94a2f4d67df544b1a5f49c4c42c6d', 'check_cell_approx__2b6911a42a0fb61cba672fa8735595e9', 'check_cell_value_numeric__0ba689dbed2285037ec7f44756763cc8', 'check_cell_text__5eff620ab5e8861fd0f5d3b257b5ca40', 'check_cell_numeric_value__27fb6350ba8922eab313e55ef2751ee9', 'check_sheet2_avg__843226a2563502226c076bf6effbabf5', 'check_sheet2_column_data__8a87fce5953091d60c1165596b42e521', 'check_xlsx_freeze_pane__7c0a43d93f3ca3bc0e18377c9b4b22b9', 'check_cell_approx__e7bff0bae84b8431a03877cfc4fe4751', 'check_two_cells__5caee571623589595fc08cd4d0ebf084', 'check_cell_approx__3d7517c3edf35996bcb8466dea369ae2', 'check_cell_exact_value__774d1bad2b697745d27771510730ed1b', 'check_cell_value__936e4dec15bc5593ae1b9d4f9d1cf524', 'check_cell_value__269bac677b4d800ce51686bd69d7a5b3', 'check_cell_bgcolor__2562d5d5c00c43d3ebb3cc2320483fc3', 'check_cell_value__9e8bcb68d0493b0f9f821969711fa470', 'check_xlsx_zoom_and_freeze__048259c26ab013ccdf7c1dac5c5ed24c', 'check_cells_c8_c11__6709fbc098f9a1f484c1ad8d161e48b0', 'check_cell_numeric_value__0ff3d8d538e7cd587b08458b9c790695', 'check_cell_value__488a286f36e304dc5e1f93ce39ded8cc', 'check_excellent_counts__4d92f18ea2c033e230b5011ee11fad9e', 'check_cell_values__3d93b5527badb1dca77f44852a74fbe3', 'check_cell_value__b28ddd95dc90fbf8e65579083da3a568', 'check_cell_value_numeric__4abf9f9f6bbbc1d8a7c4b9c7e898d79c', 'check_filtered_sheet__c2e2ad1d381d1d27297ba9392b0d9cea', 'check_sheet_names__7b9538bff9f96ac19c0aebd38f8a5f3f', 'check_xlsx_headers__79ab7a7fafe67a8a56664d433902c4eb', 'check_sheet_names__424814124c3140f33ad68e6e88a912a8', 'check_cell_numeric_value__d5a8ae815c3da82f8081ee65f4e2cea1', 'check_xlsx_sum__b86ee128372bf8371e6a9412924e6faf', 'check_cell_numeric_close__d0a50343ffaff2e09611bc0e4f880cee', 'check_sheet2_maxmin__56726039aefd1f0c57475732eb3ab1f0', 'check_cell_value__f782061fb40a07e3f7bb2ac7aceef53a', 'check_sum_formula__9956b78e7ee895cd5df580471d77a6ec', 'check_cell_value__fec4261d3cbd8743e0d3445b3fa0111f', 'check_sheet2_layout__8c3665733cda1e13a57dc5b562c81969', 'check_cell_value__02d2303f1946e73506bf1a1871c3326e', 'check_cell_numeric_value__8ed2a3a912f306294d5ece5e575ff288', 'check_cell_text__4eed649bcee7f8431b3567000e60f20f', 'check_cell_text_value__9ff5dc914928ae745b032698ddd720c0', 'check_xlsx_zoom__9fd4567c6d7b24fec1b69292feb79e7d', 'check_cell_value__2fc57009f3993d012dcfddb9200ab322', 'check_cell_numeric_close__3aaf5eb76a2414410378f220bee2678f', 'check_cell_value__02607caca98b58ee759729206fe827f7', 'check_xlsx_initials__06e4265c097251cdade138725012bf19', 'check_cell_value__2650b8c230f7272f8553cf2f429cf59a', 'check_cell_numeric_close__12cf17c32a6915e1266996d25d563e2a', 'check_xlsx_sort__08b8e71e292866610cba11baa24dc9d9', 'check_cell_value__a79441059a79951e2976eaf683658a71', 'check_two_cells__8918eac5e4b67ba092cf5e7258979923', 'check_cell_value__54e9a68f9afb8559eff309088793c132', 'check_cell_value__3dfc90e20421740ec2b3449751e1ede4', 'check_cell_contains__d4eda5c6ae488cf666c27c91b5fbd879', 'check_cell_value__bbbcaf548e27f0a1a85c1f51c758e3f5', 'check_cell_value__ab8a45e026223c32b7b22c741f1bc3ca', 'check_cell_text__52b5a1c85fa3a63d940f3aeb966674a8', 'check_cell_numeric_value__11373fb010a9f0df0f8a282c8597081e', 'check_first_sheet__1a11ade1a8ab5d9c5b9f1e6b66e05fc9', 'check_cell_value_numeric__da67ff83a4f2f562e6cd5555f90c78ad', 'check_cell_value__da4a3fddcbcd3738d01b04b1ba353fc4', 'check_cell_numeric_value__d3f5a4ed5eae2c51d540e2a10864dc43', 'check_cell_values__456143f49bf34353a5e90a3c287b15a6', 'check_sheet2_merge_pattern__3f03b2d5555a7e222477184511a0b152', 'check_xlsx_clevel_count__802aa06dfc0f94fbf79660d5d540798f', 'check_xlsx_sorted_desc__db00ed878507fde13aefe6b71212bcbd', 'check_cell_value__8afea61a57c3ecda7ab6098e540526c4', 'check_sheet2_merge_text__221a8bb1dbb77d2fdb6346ea5a7576e4', 'check_cell_string_value__04a818774a969e111c0b34c8156dff9c', 'check_cell_bgcolor__fabb4ab7d265226ac33712a7dcf2e535', 'check_cell_value_numeric__85f75e16b65eecbecc9a6c1bce611233', 'check_xlsx_email_cells__93729a4c1a924bf8a2b603d4f55d7635', 'check_cell_value__cb33c763b3ef0ae13364585aa1e6bfa3', 'check_filtered_sheet__7192ac900176bf2dc148847f36558198', 'check_cell_value__3f903c88c3cac4cb7f6f63b58e6790e4', 'check_sheet2_month_avg__d92436572f728d9c9f2cb1a9c19eb76d', 'check_cells_b3_b6__f927596325f5c4e25b12a463bfafbace', 'check_xls_cell_value__2674782c2441965e8475070c5235e39b', 'check_header_cells__25536421a644aed93ef4482447c70e86', 'check_pptx_table_cell__1546f90c5e4d91a5767b01e5f0a56119', 'check_cell_numeric__13d011eac188a961b7fc2a44b3a2069f', 'check_cell_numeric_close__7aae35559065c2f6c82a9e66b7d55095', 'check_cell_value__cb5b990d93f87d0f04a6150ee1d23a05', 'check_multi_cell_values__7424fba2d2a12056c33da3db49551203', 'check_cell_bgcolor__0e55cf8187f0e07889b1109cdd0b266f', 'check_cell_value__990c97f8e62011f017525e7c1f376ba9', 'check_xlsx_qty_sort__d422460342c18f4ab4e35ab11ffce7d8', 'check_cell_value__949a54a656a59ff2d23607cb949dacfc', 'check_cell_value_numeric__28ae1ff6307904e177d5ad2de88ea142', 'check_sheet_names__a164952ae6b41142a59183faf6bedadf_qw35sft2_d3debb98', 'check_merged_cell_value__48a88b5bb3b362cc9c78cc671cedfc7c_qw35sft2_7a9956ec', 'check_cell_numeric_value__0762adbae4ad235c185f47096cf64c91_qw35sft2_0a2abca9', 'check_sheet2_state__68824c07a6606ac16add37bf4765401b_qw35sft2_5cf70407', 'check_seqno_and_sheet_name__83227706efd58da00bb8213cee393c7a_qw35sft2_523e02bd', 'check_cell_value__cfba55e4566ed0373f9b47631b661201_qw35sft2_79ef5ec5', 'check_len_formula__29aa3f16a2f40add619bb0073f57ed2e_qw35sft2_55c5c202', 'check_cell_value_equals__d4b207925bd2d0c3ababe1e319fa8b1f_qw35sft2_5621d4f6', 'check_calc_earned__73e90374b4fd7034756570a58d380c35_qw35sft2_61998100', 'check_na_count_in_cell__1ed2b6f96cb9e3177c42c0765213bb8c_qw35sft2_65832f34', 'check_multi_cell_grade__3b5b28cd38bc58813bd97e9b2f5d4ad4_qw35sft2_e971eac3', 'check_avg_age_cell__c483fe0106f1e0b2430ad19cd97c77c0_qw35sft2_18c5427e', 'check_calc_sheet1_net_income__6142d421fefd6b784b8ad81070a58350_qw35sft2_9fa35d28', 'check_sheet2_sorted_revenue__ff6f4e09739d1959467ff60470ffe2bd_qw35sft2_dae27611', 'check_cell_string__9ab5f66fb165e5668bd5b38cac0c73c2_qw35sft2_0f40b827', 'check_cells_red__03a99dd86c6d30983d2467bc1177489c_qw35sft2_ef4d5f4f', 'check_xlsx_transposed_cells__f7593e4503a45a853c1ae96cb08aaf92_qw35sft2_b1c16c4d', 'check_calc_sheet_renamed__ecf996871673b767fb6d252021f29e29_qw35sft2_cf744cd5', 'check_xlsx_zone2_row_totals__6ef64aee01b9e4fef352e390ba82e0ce_qw35sft2_d34b1ef4', 'check_sheet2_three_columns__74fb6bc2ea3e8707d0a1b3dd5202c02d_qw35sft2_4ffd667d', 'check_ramp_accel_cells__f874da4b0492e2e023290240d31f1c8d_qw35sft2_083dc380', 'check_sheet2_pct_format__6c9c981a9ca6d475ebfb8637fc973b0b_qw35sft2_c08948ee', 'check_xlsx_padded_max__7fad2dd93f80fb906237dbbf8bdc9fb2_qw35sft2_527a04cd', 'check_xlsx_spent_and_date_format__2b619f85a4e6da743c7b9581de6419ad_qw35sft2_99cbf667', 'check_calc_pivot_and_sort__6bb297a00c63757f69bb3bc219190d5b_qw35sft2_5e630c32', 'check_cell_string__13e2a9cff490c6125c903160099b9b7a_qw35sft2_d6b97dfc', 'check_month_pivot_sheet2__360bd21f4036eac9068c6b83aba7e2e2_qw35sft2_a4406367', 'check_freeze_and_sheet__c058d8a3342bf3e7968eb49c80bec94d_qw35sft2_c4c06d2f', 'check_cell_numeric__e7e54b6e523b2a67e0ff77298aeb57dc_qw35sft2_c16bee5a', 'check_cell_numeric__57a509452fa7aaef388db657e08eeb87_qw35sft2_462d120f', 'check_cell_value_equals__68e5aa381affaca2d81c1d380acacc93_qw35sft2_c346be8b', 'check_cell_value__009b0a2429b483dc78b74a0168744c1a_qw35sft2_f16b38bd', 'check_multi_cell_text__142de785ea8e1e7c66581a5013a34ff3_qw35sft2_99cabfcf', 'check_sheet_names_partial__051a1a86398cf9c50b041fec4af58dd3_qw35sft2_2917bf35', 'check_cell_grade__b0dcd11a1bdcc1e13016c1ab56479e5f_qw35sft2_0d0437cb', 'check_cell_string__76b216d4ab5877594f0cff145afca3ba_qw35sft2_75ae6314', 'check_calc_summary_sheet__878e408cd8149dcd33226991d9a64c87_qw35sft2_b8e83f56', 'check_sheet1_net_income_sheet2_revenue__5d938927d392c1d220c2599bb703f167_qw35sft2_63e26994', 'check_calc_total_sales__dc5d11c569bf87971ba13d9e911e3365_qw35sft2_1d6bedbb', 'check_calc_earned_and_tax__2412d7535e20b0348d2c25accb946746_qw35sft2_41466692', 'check_cells_red__1a5d60ee39c11d25031d40117b38947f_qw35sft2_b467b721', 'check_sheet_renamed__0c20409273ae865247912386f480b246_qw35sft2_73a65b87', 'check_xlsx_transposed_sorted__f581e5220b9d6fb8293dc4e0f669a755_qw35sft2_1d570f56', 'check_xlsx_zone1_row_totals__3aa70d99f23178c4e0eb339ad1e35c2e_qw35sft2_9d875225', 'check_sheet2_with_grand_total__0378b78b340823f75ef2df5ebaa9b121_qw35sft2_947a6f38', 'check_sheet2_sorted__2fc0837f7c43e49c220e69070cca644b_qw35sft2_628d001d', 'check_xlsx_padded_count__50235d7deccc7f57d815dc66d0e983bc_qw35sft2_3b6a4b8c', 'xlsx_spent_format_total_row__40a8f46c1e43b9d32b31b2230ffda8db_qw35sft2_4e899856', 'check_sheet_rename_and_pdf__57d2609933a3290edf4943a154f0921f_qw35sft2_c493ce14', 'check_calc_pivot_and_total__2c779926d5355737435372c9e97f0d18_qw35sft2_cdf14c34', 'check_pivot_and_sorted__3da7df8645c1252be37e4e927311c6ef_qw35sft2_6755b7e0', 'check_cell_approx__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_4b449af3', 'check_pivot_and_sum__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_700872ff', 'check_sheet_exists__0ed76c34ce42b49c8d0f69a316afbc19_qw35sft2_e4d74bcc', 'check_cell_text__fd1c3b4eab4669bd9848fe12e4f0d1d4_qw35sft2_9cce5d4b', 'check_cells_ordered_values__1d24482d98b418ee12818e4d70230d5c_qw35sft2_5772c294', 'check_sheet_exists__08c531e403541f6bab4f59a478d0e6c2_qw35sft2_69a71ad6', 'check_sheet_and_cell__ebd5692a1e5abf6c2058f3bfb502b3de_qw35sft2_24eddc65', 'check_calc_sheet2_3col__581360a1329c03f025e39497ca2b0766_qw35sft2_1a12ad34', 'check_calc_earned_and_hours__ddba4585b2adadf3e6cdf7efce56b822_qw35sft2_e04c2448', 'check_cells_red__cc3c6531bd3eb1e9b2189f82dea73399_qw35sft2_e0fd9222', 'check_cell_value_equals__9d8210ef208924a20f5946d708d3b9b8_qw35sft2_a6259fb1', 'check_pptx_table_cells__d704f09ae59c96034fafb3c1ef256369_qw35sft2_e6a9ffc5', 'check_table_7x5_last_cell__818c616cbd97d0d5bf2ff81918416501_qw35sft2_1a57fa14', 'check_table_7x5_with_cell__2aa481b876a15405aedb5c91a8fc78e9_qw35sft2_8d4c3375', 'check_xlsx_cell_val__2feaa23241a59b8898b31056b9df1f85_qw35sft2_706eb125', 'check_file_numeric_content__f36e167c6b956664e01e5e47712b671b_qw35sft2_4e4217aa', 'check_xlsx_freeze_panes__b12f1d02d5521f031b94e1747a1c556a_qw35sft2_3f46aafb', 'check_xlsx_year_column__c99ee3ea1307143970f7d9489f51bcb4_qw35sft2_457e30c1', 'check_xlsx_cell_value__5777a0629a4e670ae915b1fe64e58378_qw35sft2_e7be0bcf', 'check_xlsx_cell_val__61d78e261ca58a364d140d5219ceb5c7_qw35sft2_922f8d42', 'check_xlsx_header_align__7d3ab32fb77b1ea64cfabfdcd56b69aa_qw35sft2_c23af72a', 'check_xlsx_row_data__5abc6535b144407839fbbe3d8a497678_qw35sft2_41d3f63e', 'check_xlsx_cell_val__9d1e0f5a79d29ecfac0c206e85fa98e9_qw35sft2_73735a2c', 'check_unseen_sheet_headers__ac461f3f5169eddd647d2f7d2aad85a2_qw35sft2_0281101b', 'check_xlsx_cell_value__9d415d4db27ab3e1c21d0be924167f0e_qw35sft2_fc10d580']

def check_cell_value__3a7517792dc1ec4e3be4f857bdb0950f(result, expected, **options):
    """Check if cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        tolerance = expected.get('tolerance', 0.01)
        if abs(actual_num - expected_num) <= abs(expected_num) * tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        pass
    return 1.0 if str(actual).strip() == str(expected_val).strip() else 0.0

def check_xlsx_reversed_names__d8dd1c41700e60d858e6fd9c2f1f5451(result, expected, **options):
    """Check that full names in 'LastName, FirstName' format are correct.

    Partial credit:
    - 0.2: Column E has a 'Full Name' header
    - 0.3: First name and last name columns are populated
    - 0.5: Full names match expected 'LastName, FirstName' format
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_names = expected.get('expected_full_names', [])
    header = result.get('header_e')
    if header and ('name' in str(header).lower() or 'full' in str(header).lower()):
        score += 0.2
    first_names = result.get('first_names', [])
    last_names = result.get('last_names', [])
    non_null_fn = sum((1 for v in first_names if v is not None))
    non_null_ln = sum((1 for v in last_names if v is not None))
    if non_null_fn >= 20 and non_null_ln >= 20:
        score += 0.3
    actual_names = result.get('full_names', [])
    if len(actual_names) == len(expected_names) and len(expected_names) > 0:
        matches = 0
        for (actual, exp) in zip(actual_names, expected_names):
            if actual is not None and str(actual).strip() == str(exp).strip():
                matches += 1
        ratio = matches / len(expected_names)
        score += 0.5 * ratio
    return min(score, 1.0)

def check_cell_approx__3440d9b31b1654ffe8da4da80be4e5fc(result, expected, **options):
    """Check if cell value approximately matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    try:
        actual_f = float(actual)
    except (TypeError, ValueError):
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    if expected_val is None:
        return 0.0
    if abs(actual_f - float(expected_val)) <= tolerance:
        return 1.0
    return 0.0

def check_cell_numeric_value__db45e57fbcc1253e5b0f6029b5341c3e(result, expected, **options):
    """Check if a cell's numeric value matches expected within tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_ods_to_xlsx__d8b4b0a7ce586e5e25f592627cc186b5(result, expected, **options):
    """Check ODS to XLSX conversion: terminal usage + valid XLSX output."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    score = 0.0
    if result.get('used_terminal') == 'use terminal':
        score += 0.5
    expected_rows = expected.get('expected_min_rows', 5000)
    if result.get('file_exists'):
        row_count = result.get('row_count', 0)
        if row_count >= expected_rows and result.get('has_data'):
            score += 0.5
        elif row_count > 0 and result.get('has_data'):
            score += 0.25
    return min(score, 1.0)

def check_cell_numeric__f3f19b96be4e433759cb0f99e36be750(result, expected, **options):
    """Check if a cell has the expected numeric value with tolerance."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_cell_b3_value__0698f7f5da73d25e009ccfe66a594e11(result, expected, **options):
    """Check if cell B3 has the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    if str(actual).strip().lower() == str(expected_val).strip().lower():
        return 1.0
    return 0.0

def check_cell_value__4a0f3185fd3c30bf019ccadcf5573a11(result, expected, **options):
    """Check if cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        tolerance = expected.get('tolerance', 0.01)
        if abs(actual_num - expected_num) <= abs(expected_num) * tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        pass
    return 1.0 if str(actual).strip() == str(expected_val).strip() else 0.0

def check_multi_cells__677226b05224be2ecbd2de2e693f9c8e(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not expected_values:
        return 0.0
    total = len(expected_values)
    matched = 0
    for (cell, exp_val) in expected_values.items():
        act_val = actual_values.get(cell)
        if act_val is None:
            continue
        try:
            act_num = float(act_val)
            exp_num = float(exp_val)
            if abs(act_num - exp_num) < 0.5:
                matched += 1
                continue
        except (TypeError, ValueError):
            pass
        act_str = str(act_val).strip().rstrip('%').strip()
        exp_str = str(exp_val).strip().rstrip('%').strip()
        try:
            if abs(float(act_str) - float(exp_str)) < 0.5:
                matched += 1
                continue
        except (TypeError, ValueError):
            pass
        if str(act_val).strip().lower() == str(exp_val).strip().lower():
            matched += 1
    return matched / total if total > 0 else 0.0

def check_cell_sum_value__dc007ee2aea5c445c94303ee0d01fcf4(result, expected, **options):
    """Check if cell contains the expected sum value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None or expected_value is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_value)) < 0.01:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_filtered_sheet__0d0f959234d0c494b6d7c17cb5158836(result, expected, **options):
    """Check that a filtered sheet contains the correct movies with proper sorting.

    Partial credit:
    - 0.25: Sheet exists with correct name
    - 0.25: Headers match original format
    - 0.25: Correct number of rows (within tolerance)
    - 0.25: Data is correctly sorted
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_sheet = expected.get('expected_sheet_name', '')
    if result.get('sheet_name', '').strip().lower() == expected_sheet.strip().lower():
        score += 0.25
    expected_headers = expected.get('expected_headers', [])
    actual_headers = result.get('headers', [])
    if len(actual_headers) >= len(expected_headers):
        match = all((str(a).strip().lower() == str(e).strip().lower() for (a, e) in zip(actual_headers, expected_headers)))
        if match:
            score += 0.25
    expected_count = expected.get('expected_row_count', 0)
    actual_count = result.get('row_count', 0)
    if actual_count == expected_count:
        score += 0.25
    elif abs(actual_count - expected_count) <= 2:
        score += 0.125
    sort_key = expected.get('sort_key', '')
    sort_order = expected.get('sort_order', 'asc')
    rows = result.get('rows', [])
    if rows and sort_key:
        values = []
        for r in rows:
            v = r.get(sort_key)
            if v is not None:
                try:
                    values.append(float(str(v)))
                except (ValueError, TypeError):
                    values.append(0)
        if len(values) >= 2:
            if sort_order == 'desc':
                is_sorted = all((values[i] >= values[i + 1] for i in range(len(values) - 1)))
            else:
                is_sorted = all((values[i] <= values[i + 1] for i in range(len(values) - 1)))
            if is_sorted:
                score += 0.25
    return min(score, 1.0)

def check_xlsx_new_row__246ef0821ed926b9e34f1369cfdfe795(result, expected, **options):
    """Check if a new row was added with expected values. Partial credit per column."""
    if not result or isinstance(result, str) or result.get('error'):
        return 0.0
    row_values = result.get('row_values', {})
    expected_values = expected.get('expected_row', {})
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (col, exp_val) in expected_values.items():
        actual = row_values.get(col)
        if actual is not None and str(actual).strip().lower() == str(exp_val).strip().lower():
            correct += 1
    return correct / total if total > 0 else 0.0

def check_range_cells__980775095293f828bec4466da30db27b(result, expected, **options):
    """Check if cell values in a range match expected values. Supports partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    cells = result.get('cells', {})
    expected_cells = expected.get('expected_cells', {})
    if not expected_cells:
        return 0.0
    total = len(expected_cells)
    matched = 0
    for (cell_ref, expected_val) in expected_cells.items():
        actual_val = cells.get(cell_ref)
        if actual_val is not None and expected_val is not None:
            if actual_val.strip() == str(expected_val).strip():
                matched += 1
    return matched / total if total > 0 else 0.0

def check_multi_cell_values__6d0d805c4a52313050aa161e6524aaf8(result, expected, **options):
    """Check multiple cell values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual_values = result.get('values', {})
    expected_checks = expected.get('checks', {})
    if not expected_checks:
        return 0.0
    total_score = 0.0
    num_checks = len(expected_checks)
    for (cell, exp_val) in expected_checks.items():
        actual = actual_values.get(cell)
        if actual is None:
            continue
        if isinstance(exp_val, (int, float)):
            try:
                actual_num = float(actual)
                tolerance = expected.get('tolerance', 0.5)
                if abs(actual_num - float(exp_val)) <= tolerance:
                    total_score += 1.0 / num_checks
            except (TypeError, ValueError):
                pass
        elif str(actual).strip() == str(exp_val).strip():
            total_score += 1.0 / num_checks
    return min(total_score, 1.0)

def check_cell_values_v0__c8482c50afd0db1eeb95cd8ec431bb1e(result, expected, **options):
    """Check if A10 has the label and B10 has the correct employee count."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    a10 = result.get('a10', '')
    expected_label = expected.get('a10_label', 'Total Employees')
    if isinstance(a10, str) and expected_label.lower() in a10.lower():
        score += 0.3
    b10 = result.get('b10')
    expected_count = expected.get('b10_value', 7)
    if b10 is not None:
        try:
            if float(b10) == float(expected_count):
                score += 0.7
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_xlsx_cell_value__fe91ae8d0e7c34f2eef1f4db512bed2a(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None and expected_value is None:
        return 1.0
    if actual is None or expected_value is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    if actual_str.lower() == expected_str.lower():
        return 0.8
    return 0.0

def check_sheet_names__9889ac58b2fb5d2c955b2821b86f86b7(result, expected, **options):
    """Check if sheet names match expected names exactly (names, order, count)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_names = result.get('sheet_names', [])
    expected_names = expected.get('expected_sheet_names', [])
    if not expected_names:
        return 0.0
    if actual_names == expected_names:
        return 1.0
    score = 0.0
    per_name = 1.0 / len(expected_names)
    for exp_name in expected_names:
        if exp_name in actual_names:
            score += per_name
    extra_count = len(actual_names) - len(expected_names)
    if extra_count > 0:
        score -= 0.3 * extra_count
    return max(min(score, 0.99), 0.0)

def check_multi_cell_values__8a05426d089eda14388178863fc78e0b(result, expected, **options):
    """Check multiple cell values with partial credit and tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    total_weight = sum((c.get('weight', 1.0) for c in checks))
    score = 0.0
    for check in checks:
        cell = check.get('cell')
        expected_val = check.get('expected_value')
        check_type = check.get('check_type', 'numeric')
        tolerance = check.get('tolerance', 0.01)
        weight = check.get('weight', 1.0)
        actual = result.get(cell)
        if actual is None:
            continue
        if check_type == 'numeric':
            try:
                if abs(float(actual) - float(expected_val)) <= tolerance:
                    score += weight
            except (TypeError, ValueError):
                pass
        elif check_type == 'string':
            if str(actual).strip().lower() == str(expected_val).strip().lower():
                score += weight
    return min(score / total_weight, 1.0)

def check_sheet_name__e743a4ede9bcdbd814ac13f073acba4a(result, expected, **options):
    """Check if worksheet name matches expected name."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_name = result.get('sheet_name', '')
    expected_name = expected.get('expected_name', '')
    if not actual_name or not expected_name:
        return 0.0
    return 1.0 if actual_name.strip() == expected_name.strip() else 0.0

def check_sheet_renamed__64efe2d9a5895dad59edc0eb16307294(result, expected, **options):
    """Check if a sheet was renamed: old name gone, new name present."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    sheet_names = result.get('sheet_names', [])
    old_name = expected.get('old_name', '')
    new_name = expected.get('new_name', '')
    score = 0.0
    if new_name in sheet_names:
        score += 0.5
    if old_name not in sheet_names:
        score += 0.5
    return score

def check_multi_cell__1393f3a73542b23a8e048e3ef0292819(result, expected, **options):
    """Check multiple cell values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', {})
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    weight = 1.0 / len(checks)
    for (cell, expected_val) in checks.items():
        actual = values.get(cell)
        if actual is None and expected_val is None:
            score += weight
        elif actual is not None and expected_val is not None:
            actual_str = str(actual).strip().lower()
            expected_str = str(expected_val).strip().lower()
            if actual_str == expected_str:
                score += weight
    return min(score, 1.0)

def check_cell_value__3ad3a7a9e15d5bbcab97b429c7a2d70d(result, expected, **options):
    """Check that cell H2 contains the expected SUMIF value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_xlsx_cell_value__5729380214ac80099631e898f8c00bcc(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None and expected_value is None:
        return 1.0
    if actual is None or expected_value is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    if actual_str.lower() == expected_str.lower():
        return 0.8
    return 0.0

def check_cell_numeric_value__c13e5303d821a816dfe6dbe089834e52(result, expected, **options):
    """Check if cell value matches expected numeric value with tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_cell_value_numeric__ad5647b3d1b47a9e1996a6e69ca9562d(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.5)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
    except (TypeError, ValueError):
        return 0.0
    if abs(actual_num - expected_num) <= tolerance:
        return 1.0
    return 0.0

def check_cell_numeric_value__67a1f86b27a24adb60e8a183c35b0852(result, expected, **options):
    """Check if a cell contains the expected numeric value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_cells_replaced__1cd51ca9c28aca9840aa590661974218(result, expected, **options):
    """Check if specified cells have been replaced with expected value. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_cells = expected.get('expected_cells', {})
    if not expected_cells:
        return 0.0
    correct = 0
    total = len(expected_cells)
    for (cell_ref, expected_val) in expected_cells.items():
        actual_val = values.get(cell_ref)
        if actual_val is not None and str(actual_val).strip().lower() == str(expected_val).strip().lower():
            correct += 1
    return correct / total if total > 0 else 0.0

def check_xlsx_filenames__aad13293a214d4a4614be1039d459c78(result, expected, **options):
    """Check if xlsx column contains expected values with partial credit.
    Each correct value in correct position gets equal credit.
    """
    if result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    score = 0.0
    per_item = 1.0 / len(expected_values)
    for (i, exp) in enumerate(expected_values):
        if i < len(actual_values):
            actual_clean = actual_values[i].strip().lower()
            exp_clean = exp.strip().lower()
            if exp_clean == actual_clean or exp_clean in actual_clean:
                score += per_item
    return min(score, 1.0)

def check_cell_string_value__69f94a2f4d67df544b1a5f49c4c42c6d(result, expected, **options):
    """Check if a cell's string value matches expected."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    if expected_val is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_val).strip()
    if actual_str == expected_str:
        return 1.0
    if actual_str.lower() == expected_str.lower():
        return 0.8
    return 0.0

def check_cell_approx__2b6911a42a0fb61cba672fa8735595e9(result, expected, **options):
    """Check if cell value approximately matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    try:
        actual_f = float(actual)
    except (TypeError, ValueError):
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    if expected_val is None:
        return 0.0
    if abs(actual_f - float(expected_val)) <= tolerance:
        return 1.0
    return 0.0

def check_cell_value_numeric__0ba689dbed2285037ec7f44756763cc8(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 1.0)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
    except (TypeError, ValueError):
        return 0.0
    if abs(actual_num - expected_num) <= tolerance:
        return 1.0
    return 0.0

def check_cell_text__5eff620ab5e8861fd0f5d3b257b5ca40(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value', '')
    expected_value = expected.get('expected_value', '')
    if actual is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    if expected_str.lower() in actual_str.lower():
        return 0.5
    return 0.0

def check_cell_numeric_value__27fb6350ba8922eab313e55ef2751ee9(result, expected, **options):
    """Check if cell value matches expected numeric value with tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_sheet2_avg__843226a2563502226c076bf6effbabf5(result, expected, **options):
    """Check average Revenue and Expenses on Sheet2."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header_a = str(result.get('header_a', '')).lower()
    if 'average' in header_a and 'revenue' in header_a:
        score += 0.15
    header_b = str(result.get('header_b', '')).lower()
    if 'average' in header_b and 'expense' in header_b:
        score += 0.15
    try:
        actual_a = float(result.get('value_a', 0))
        expected_a = float(expected.get('expected_avg_revenue', 0))
        if expected_a > 0 and abs(actual_a - expected_a) / expected_a < 0.01:
            score += 0.35
    except (TypeError, ValueError):
        pass
    try:
        actual_b = float(result.get('value_b', 0))
        expected_b = float(expected.get('expected_avg_expenses', 0))
        if expected_b > 0 and abs(actual_b - expected_b) / expected_b < 0.01:
            score += 0.35
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_sheet2_column_data__8a87fce5953091d60c1165596b42e521(result, expected, **options):
    """Check if Sheet2 contains the expected column data with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_header = expected.get('expected_header')
    if result.get('header') == expected_header:
        score += 0.2
    expected_values = expected.get('expected_values', [])
    actual_values = result.get('values', [])
    if expected_values and actual_values:
        matches = 0
        total = len(expected_values)
        for i in range(min(len(actual_values), total)):
            if actual_values[i] == expected_values[i]:
                matches += 1
        if total > 0:
            score += 0.8 * (matches / total)
    return min(score, 1.0)

def check_xlsx_freeze_pane__7c0a43d93f3ca3bc0e18377c9b4b22b9(result, expected, **options):
    """Check if freeze pane is set to the expected cell reference."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_freeze = result.get('freeze_panes')
    expected_freeze = expected.get('expected_freeze')
    if actual_freeze is None and expected_freeze is None:
        return 1.0
    if actual_freeze == expected_freeze:
        return 1.0
    return 0.0

def check_cell_approx__e7bff0bae84b8431a03877cfc4fe4751(result, expected, **options):
    """Check that a cell value matches expected with tolerance.

    Expected keys: expected_total (float), tolerance (float, default 0.01).
    Also checks for a "Total" label row.
    Scoring: 0.5 for Total label existing, 0.5 for correct sum value.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    exp_total = float(expected.get('expected_total', 0))
    tolerance = float(expected.get('tolerance', 0.02))
    score = 0.0
    total_rows = result.get('total_rows', [])
    if total_rows:
        score += 0.5
        for tr in total_rows:
            amount = tr.get('amount')
            if amount is not None:
                try:
                    if abs(float(amount) - exp_total) < tolerance:
                        score += 0.5
                        break
                except (ValueError, TypeError):
                    pass
    else:
        c7_val = result.get('C7')
        if c7_val is not None:
            try:
                if abs(float(c7_val) - exp_total) < tolerance:
                    score += 1.0
            except (ValueError, TypeError):
                pass
    return min(score, 1.0)

def check_two_cells__5caee571623589595fc08cd4d0ebf084(result, expected, **options):
    """Check two cell values with partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_val1 = expected.get('expected_cell1')
    expected_val2 = expected.get('expected_cell2')
    actual_val1 = result.get('cell1_value')
    actual_val2 = result.get('cell2_value')
    if actual_val1 is not None and str(actual_val1).strip().lower() == str(expected_val1).strip().lower():
        score += 0.5
    if actual_val2 is not None and expected_val2 is not None:
        try:
            if abs(float(actual_val2) - float(expected_val2)) < 0.01:
                score += 0.5
        except (ValueError, TypeError):
            if str(actual_val2).strip() == str(expected_val2).strip():
                score += 0.5
    return min(score, 1.0)

def check_cell_approx__3d7517c3edf35996bcb8466dea369ae2(result, expected, **options):
    """Check if cell value approximately matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    try:
        actual_f = float(actual)
    except (TypeError, ValueError):
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    if expected_val is None:
        return 0.0
    if abs(actual_f - float(expected_val)) <= tolerance:
        return 1.0
    return 0.0

def check_cell_exact_value__774d1bad2b697745d27771510730ed1b(result, expected, **options):
    """Check if a cell value matches the expected value exactly (with numeric type tolerance)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) < 0.001:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_cell_value__936e4dec15bc5593ae1b9d4f9d1cf524(result, expected, **options):
    """Check if a cell value matches the expected value (case-insensitive string comparison)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    if isinstance(expected_val, str) and isinstance(actual, str):
        if actual.strip().lower() == expected_val.strip().lower():
            return 1.0
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual) == str(expected_val):
        return 1.0
    return 0.0

def check_cell_value__269bac677b4d800ce51686bd69d7a5b3(result, expected, **options):
    """Check if a cell value matches the expected numeric value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_val):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_cell_bgcolor__2562d5d5c00c43d3ebb3cc2320483fc3(result, expected, **options):
    """Check if specified cells have the expected background color.
    Supports partial credit: score = fraction of cells with correct color.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    colors = result.get('colors', {})
    expected_color = expected.get('expected_color', '').upper()
    expected_cells = expected.get('cells', [])
    if not expected_cells or not expected_color:
        return 0.0
    correct = 0
    for cell_ref in expected_cells:
        actual = colors.get(cell_ref)
        if actual and expected_color in actual.upper():
            correct += 1
    return correct / len(expected_cells) if expected_cells else 0.0

def check_cell_value__9e8bcb68d0493b0f9f821969711fa470(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None or expected_value is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_value)) < 0.01:
            return 1.0
    except (ValueError, TypeError):
        if str(actual) == str(expected_value):
            return 1.0
    return 0.0

def check_xlsx_zoom_and_freeze__048259c26ab013ccdf7c1dac5c5ed24c(result, expected, **options):
    """Check both zoom level and freeze pane with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_zoom = result.get('zoom')
    expected_zoom = expected.get('expected_zoom')
    if actual_zoom is not None and expected_zoom is not None and (actual_zoom == expected_zoom):
        score += 0.5
    actual_freeze = result.get('freeze_panes')
    expected_freeze = expected.get('expected_freeze')
    if actual_freeze == expected_freeze:
        score += 0.5
    return min(score, 1.0)

def check_cells_c8_c11__6709fbc098f9a1f484c1ad8d161e48b0(result, expected, **options):
    """Check if cells C8-C11 all have the expected value, with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_val = expected.get('expected_value')
    cells = expected.get('cells', ['C8', 'C9', 'C10', 'C11'])
    if not expected_val or not values:
        return 0.0
    score = 0.0
    per_cell = 1.0 / len(cells)
    for cell in cells:
        actual = values.get(cell)
        if actual is not None and str(actual).strip().lower() == str(expected_val).strip().lower():
            score += per_cell
    return min(score, 1.0)

def check_cell_numeric_value__0ff3d8d538e7cd587b08458b9c790695(result, expected, **options):
    """Check if a cell's numeric value matches expected within tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_cell_value__488a286f36e304dc5e1f93ce39ded8cc(result, expected, **options):
    """Check if a single cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.5:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_excellent_counts__4d92f18ea2c033e230b5011ee11fad9e(result, expected, **options):
    """Check header in AA1 and COUNTIF values for Excellent ratings in AA2:AA8."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header = result.get('header', '')
    expected_header = expected.get('header', 'Excellent Count')
    if isinstance(header, str) and expected_header.lower() in header.lower():
        score += 0.2
    counts = result.get('counts', [])
    expected_counts = expected.get('counts', [])
    if counts and expected_counts:
        per_employee = 0.8 / len(expected_counts)
        for i in range(min(len(counts), len(expected_counts))):
            if counts[i] is not None:
                try:
                    if int(float(counts[i])) == expected_counts[i]:
                        score += per_employee
                except (ValueError, TypeError):
                    pass
    return min(score, 1.0)

def check_cell_values__3d93b5527badb1dca77f44852a74fbe3(result, expected, **options):
    """Check cell values with type-flexible comparison. Supports partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', {})
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    weight_per_check = 1.0 / len(checks)
    for (cell_ref, exp_info) in checks.items():
        actual = values.get(cell_ref)
        exp_val = exp_info.get('value')
        match_type = exp_info.get('match', 'exact')
        if match_type == 'exact':
            if actual is not None and str(actual).strip() == str(exp_val).strip():
                score += weight_per_check
        elif match_type == 'numeric':
            try:
                if abs(float(actual) - float(exp_val)) < 0.01:
                    score += weight_per_check
            except (TypeError, ValueError):
                pass
    return min(score, 1.0)

def check_cell_value__b28ddd95dc90fbf8e65579083da3a568(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_value):
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_cell_value_numeric__4abf9f9f6bbbc1d8a7c4b9c7e898d79c(result, expected, **options):
    """Check if a cell value matches expected numeric value with tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_filtered_sheet__c2e2ad1d381d1d27297ba9392b0d9cea(result, expected, **options):
    """Check that a filtered sheet contains the correct movies with proper sorting.

    Partial credit:
    - 0.25: Sheet exists with correct name
    - 0.25: Headers match original format
    - 0.125: Correct number of rows (within tolerance)
    - 0.125: All rows have release year within expected range [1990, 1999]
    - 0.25: Data is correctly sorted
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_sheet = expected.get('expected_sheet_name', '')
    if result.get('sheet_name', '').strip().lower() == expected_sheet.strip().lower():
        score += 0.25
    expected_headers = expected.get('expected_headers', [])
    actual_headers = result.get('headers', [])
    if len(actual_headers) >= len(expected_headers):
        match = all((str(a).strip().lower() == str(e).strip().lower() for (a, e) in zip(actual_headers, expected_headers)))
        if match:
            score += 0.25
    expected_count = expected.get('expected_row_count', 0)
    actual_count = result.get('row_count', 0)
    if actual_count == expected_count:
        score += 0.125
    elif abs(actual_count - expected_count) <= 2:
        score += 0.0625
    year_key = expected.get('year_key', 'release year')
    year_min = expected.get('year_range_min', 1990)
    year_max = expected.get('year_range_max', 1999)
    rows = result.get('rows', [])
    if rows:
        all_in_range = True
        checked = 0
        for r in rows:
            v = r.get(year_key)
            if v is not None:
                try:
                    year_val = int(float(str(v)))
                    if year_val < year_min or year_val > year_max:
                        all_in_range = False
                        break
                    checked += 1
                except (ValueError, TypeError):
                    all_in_range = False
                    break
            else:
                all_in_range = False
                break
        if all_in_range and checked > 0:
            score += 0.125
    sort_key = expected.get('sort_key', '')
    sort_order = expected.get('sort_order', 'desc')
    if rows and sort_key:
        values = []
        for r in rows:
            v = r.get(sort_key)
            if v is not None:
                try:
                    values.append(float(str(v)))
                except (ValueError, TypeError):
                    values.append(0)
        if len(values) >= 2:
            if sort_order == 'desc':
                is_sorted = all((values[i] >= values[i + 1] for i in range(len(values) - 1)))
            else:
                is_sorted = all((values[i] <= values[i + 1] for i in range(len(values) - 1)))
            if is_sorted:
                score += 0.25
    return min(score, 1.0)

def check_sheet_names__7b9538bff9f96ac19c0aebd38f8a5f3f(result, expected, **options):
    """Check if sheet names match expected names in order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_names = result.get('sheet_names', [])
    expected_names = expected.get('expected_sheet_names', [])
    if not expected_names:
        return 0.0
    score = 0.0
    per_name = 1.0 / len(expected_names)
    for exp_name in expected_names:
        if exp_name in actual_names:
            score += per_name
    if actual_names == expected_names:
        return 1.0
    return min(score, 0.99)

def check_xlsx_headers__79ab7a7fafe67a8a56664d433902c4eb(result, expected, **options):
    """Check if xlsx headers match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not actual_values or not expected_values:
        return 0.0
    score = 0.0
    total = len(expected_values)
    for (i, exp) in enumerate(expected_values):
        if i < len(actual_values) and actual_values[i] is not None:
            if actual_values[i].lower().strip() == str(exp).lower().strip():
                score += 1.0 / total
    return min(score, 1.0)

def check_sheet_names__424814124c3140f33ad68e6e88a912a8(result, expected, **options):
    """Check if sheet names match expected names in order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_names = result.get('sheet_names', [])
    expected_names = expected.get('expected_sheet_names', [])
    if not expected_names:
        return 0.0
    score = 0.0
    per_name = 1.0 / len(expected_names)
    for exp_name in expected_names:
        if exp_name in actual_names:
            score += per_name
    if actual_names == expected_names:
        return 1.0
    return min(score, 0.99)

def check_cell_numeric_value__d5a8ae815c3da82f8081ee65f4e2cea1(result, expected, **options):
    """Check if a cell contains the expected numeric value."""
    if result.get('error'):
        return 0.0
    value = result.get('value')
    if value is None:
        return 0.0
    expected_value = expected.get('expected_value')
    if expected_value is None:
        return 0.0
    try:
        actual_num = float(value)
        expected_num = float(expected_value)
        if actual_num == expected_num:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_xlsx_sum__b86ee128372bf8371e6a9412924e6faf(result, expected, **options):
    """Check if the SUM formula was correctly added and produces expected value."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('has_formula'):
        score += 0.4
    expected_value = expected.get('expected_value', 780)
    actual_value = result.get('value')
    if actual_value is not None:
        try:
            if abs(float(actual_value) - float(expected_value)) < 0.01:
                score += 0.6
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_cell_numeric_close__d0a50343ffaff2e09611bc0e4f880cee(result, expected, **options):
    """Check if a cell value is numerically close to the expected value."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.1)
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_sheet2_maxmin__56726039aefd1f0c57475732eb3ab1f0(result, expected, **options):
    """Check max Revenue and min Expenses on Sheet2."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    header_a = str(result.get('header_a', '')).lower()
    if 'max' in header_a and 'revenue' in header_a:
        score += 0.15
    header_b = str(result.get('header_b', '')).lower()
    if 'min' in header_b and 'expense' in header_b:
        score += 0.15
    try:
        actual_a = float(result.get('value_a', 0))
        expected_a = float(expected.get('expected_max_revenue', 0))
        if abs(actual_a - expected_a) < 0.01:
            score += 0.35
    except (TypeError, ValueError):
        pass
    try:
        actual_b = float(result.get('value_b', 0))
        expected_b = float(expected.get('expected_min_expenses', 0))
        if abs(actual_b - expected_b) < 0.01:
            score += 0.35
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_cell_value__f782061fb40a07e3f7bb2ac7aceef53a(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_sum_formula__9956b78e7ee895cd5df580471d77a6ec(result, expected, **options):
    """Check if cell has correct SUM formula and value."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    formula = result.get('formula', '')
    if formula:
        formula_upper = str(formula).upper().replace(' ', '')
        valid_formulas = ['=SUM(C2:C8)', '=SUM(C2:C8)']
        if any((f.upper() == formula_upper for f in valid_formulas)):
            score += 0.5
        elif 'SUM' in formula_upper and 'C2' in formula_upper and ('C8' in formula_upper):
            score += 0.5
    value = result.get('value')
    if value is not None:
        try:
            actual = float(value)
            expected_val = float(expected.get('expected_value', 288.09))
            if abs(actual - expected_val) < 0.01:
                score += 0.5
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_cell_value__fec4261d3cbd8743e0d3445b3fa0111f(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_sheet2_layout__8c3665733cda1e13a57dc5b562c81969(result, expected, **options):
    """Check Sheet2 has correct title merge and row 2 headers."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    a1 = result.get('A1', {})
    if a1.get('value') == rules.get('title_value'):
        score += 0.2
    merge_cells = ['B1', 'C1', 'D1', 'E1']
    merge_count = sum((1 for c in merge_cells if result.get(c, {}).get('merged', False)))
    score += 0.3 * (merge_count / len(merge_cells))
    expected_headers = rules.get('row2_headers', {})
    header_count = 0
    total_headers = len(expected_headers)
    if total_headers > 0:
        for (cell_ref, exp_val) in expected_headers.items():
            cell_data = result.get(cell_ref, {})
            if cell_data.get('value') == exp_val:
                header_count += 1
        score += 0.5 * (header_count / total_headers)
    return min(score, 1.0)

def check_cell_value__02d2303f1946e73506bf1a1871c3326e(result, expected, **options):
    """Check if a cell value matches the expected string value (case-insensitive, trimmed)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    actual_str = str(actual).strip().lower()
    expected_str = str(expected_val).strip().lower()
    if actual_str == expected_str:
        return 1.0
    if expected_str in actual_str or actual_str in expected_str:
        return 0.5
    return 0.0

def check_cell_numeric_value__8ed2a3a912f306294d5ece5e575ff288(result, expected, **options):
    """Check if a cell value matches expected numeric value within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < tolerance:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_cell_text__4eed649bcee7f8431b3567000e60f20f(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value', '')
    expected_value = expected.get('expected_value', '')
    if actual is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    if expected_str.lower() in actual_str.lower():
        return 0.5
    return 0.0

def check_cell_text_value__9ff5dc914928ae745b032698ddd720c0(result, expected, **options):
    """Check if a cell's text value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    if str(actual).strip().lower() == str(expected_val).strip().lower():
        return 1.0
    return 0.0

def check_xlsx_zoom__9fd4567c6d7b24fec1b69292feb79e7d(result, expected, **options):
    """Check if the zoom level matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_zoom = result.get('zoom')
    expected_zoom = expected.get('expected_zoom')
    if actual_zoom is None or expected_zoom is None:
        return 0.0
    if actual_zoom == expected_zoom:
        return 1.0
    return 0.0

def check_cell_value__2fc57009f3993d012dcfddb9200ab322(result, expected, **options):
    """Check if a cell value matches the expected numeric value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_val):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_cell_numeric_close__3aaf5eb76a2414410378f220bee2678f(result, expected, **options):
    """Check if a numeric cell value is close to expected value within tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 1.0)
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_cell_value__02607caca98b58ee759729206fe827f7(result, expected, **options):
    """Check if cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        tolerance = expected.get('tolerance', 0.01)
        if abs(actual_num - expected_num) < tolerance:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_xlsx_initials__06e4265c097251cdade138725012bf19(result, expected, **options):
    """Check that initials column is correctly populated from split data.

    Partial credit:
    - 0.2: Column E has 'Initials' header
    - 0.4: First name and last name columns are populated
    - 0.4: Initials match expected (first letter of first name + first letter of last name)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_initials = expected.get('expected_initials', [])
    header = result.get('header_e')
    if header and 'nitial' in str(header).lower():
        score += 0.2
    first_names = result.get('first_names', [])
    last_names = result.get('last_names', [])
    non_null_fn = sum((1 for v in first_names if v is not None))
    non_null_ln = sum((1 for v in last_names if v is not None))
    if non_null_fn >= 20 and non_null_ln >= 20:
        score += 0.4
    actual_initials = result.get('initials', [])
    if len(actual_initials) == len(expected_initials) and len(expected_initials) > 0:
        matches = 0
        for (actual, exp) in zip(actual_initials, expected_initials):
            if actual is not None and str(actual).strip().upper() == str(exp).strip().upper():
                matches += 1
        ratio = matches / len(expected_initials)
        score += 0.4 * ratio
    return min(score, 1.0)

def check_cell_value__2650b8c230f7272f8553cf2f429cf59a(result, expected, **options):
    """Check if a cell value matches the expected numeric value."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).rstrip('0').rstrip('.') == str(expected_val).rstrip('0').rstrip('.'):
        return 1.0
    return 0.0

def check_cell_numeric_close__12cf17c32a6915e1266996d25d563e2a(result, expected, **options):
    """Check if a cell value is numerically close to expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    value = result.get('value')
    if value is None:
        return 0.0
    try:
        actual = float(value)
    except (TypeError, ValueError):
        return 0.0
    expected_val = expected.get('expected_value')
    if expected_val is None:
        return 0.0
    tolerance = expected.get('tolerance', 0.01)
    if abs(actual - float(expected_val)) <= tolerance:
        return 1.0
    return 0.0

def check_xlsx_sort__08b8e71e292866610cba11baa24dc9d9(result, expected, **options):
    """Check if spreadsheet data is sorted by column A ascending."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_sorted_ascending'):
        score += 0.7
    expected_rows = expected.get('expected_row_count', 35)
    if result.get('row_count') == expected_rows:
        score += 0.15
    first_val = result.get('first_value', '')
    expected_first = expected.get('expected_first_value', '')
    if first_val and expected_first and (expected_first in first_val):
        score += 0.15
    return min(score, 1.0)

def check_cell_value__a79441059a79951e2976eaf683658a71(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_two_cells__8918eac5e4b67ba092cf5e7258979923(result, expected, **options):
    """Check two cell values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    actual1 = result.get('cell1_value')
    expected1 = rules.get('expected_cell1')
    if actual1 is not None and str(actual1).strip() == str(expected1).strip():
        score += 0.5
    actual2 = result.get('cell2_value')
    expected2 = rules.get('expected_cell2')
    if actual2 is not None:
        try:
            if float(actual2) == float(expected2):
                score += 0.5
        except (ValueError, TypeError):
            if str(actual2).strip() == str(expected2).strip():
                score += 0.5
    return min(score, 1.0)

def check_cell_value__54e9a68f9afb8559eff309088793c132(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    try:
        if float(actual) == float(expected_val):
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_cell_value__3dfc90e20421740ec2b3449751e1ede4(result, expected, **options):
    """Check if cell value matches expected value (case-insensitive string comparison)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None and expected_val is None:
        return 1.0
    if actual is None or expected_val is None:
        return 0.0
    actual_str = str(actual).strip().lower()
    expected_str = str(expected_val).strip().lower()
    if actual_str == expected_str:
        return 1.0
    return 0.0

def check_cell_contains__d4eda5c6ae488cf666c27c91b5fbd879(result, expected, **options):
    """Check if a cell contains the expected text (case-insensitive substring match)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    value = result.get('value')
    if value is None:
        return 0.0
    actual = str(value).strip().lower()
    expected_text = str(expected.get('expected_text', '')).strip().lower()
    if not expected_text:
        return 0.0
    if expected_text in actual:
        return 1.0
    return 0.0

def check_cell_value__bbbcaf548e27f0a1a85c1f51c758e3f5(result, expected, **options):
    """Check if a cell value matches the expected numeric value."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).rstrip('0').rstrip('.') == str(expected_val).rstrip('0').rstrip('.'):
        return 1.0
    return 0.0

def check_cell_value__ab8a45e026223c32b7b22c741f1bc3ca(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_cell_text__52b5a1c85fa3a63d940f3aeb966674a8(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value', '')
    expected_value = expected.get('expected_value', '')
    if actual is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    if expected_str.lower() in actual_str.lower():
        return 0.5
    return 0.0

def check_cell_numeric_value__11373fb010a9f0df0f8a282c8597081e(result, expected, **options):
    """Check if a cell's numeric value matches expected within tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.5)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_first_sheet__1a11ade1a8ab5d9c5b9f1e6b66e05fc9(result, expected, **options):
    """Check if the first sheet name matches expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    first_sheet = result.get('first_sheet', '')
    expected_first = expected.get('expected_first_sheet', '')
    if first_sheet == expected_first:
        return 1.0
    return 0.0

def check_cell_value_numeric__da67ff83a4f2f562e6cd5555f90c78ad(result, expected, **options):
    """Check if a cell's numeric value matches expected within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else None
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_f = float(actual)
        expected_f = float(expected_val)
        if abs(actual_f - expected_f) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_cell_value__da4a3fddcbcd3738d01b04b1ba353fc4(result, expected, **options):
    """Check if cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_cell_numeric_value__d3f5a4ed5eae2c51d540e2a10864dc43(result, expected, **options):
    """Check if a cell's numeric value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_cell_values__456143f49bf34353a5e90a3c287b15a6(result, expected, **options):
    """Check cell values with type-flexible comparison. Supports partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', {})
    checks = expected.get('checks', {})
    if not checks:
        return 0.0
    score = 0.0
    weight_per_check = 1.0 / len(checks)
    for (cell_ref, exp_info) in checks.items():
        actual = values.get(cell_ref)
        exp_val = exp_info.get('value')
        match_type = exp_info.get('match', 'exact')
        if match_type == 'exact':
            if actual is not None and str(actual).strip() == str(exp_val).strip():
                score += weight_per_check
        elif match_type == 'numeric':
            try:
                if abs(float(actual) - float(exp_val)) < 1.0:
                    score += weight_per_check
            except (TypeError, ValueError):
                pass
    return min(score, 1.0)

def check_sheet2_merge_pattern__3f03b2d5555a7e222477184511a0b152(result, expected, **options):
    """Check Sheet2 has correct merge pattern and values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    checks = rules.get('checks', [])
    if not checks:
        return 0.0
    per_check = 1.0 / len(checks)
    for check in checks:
        cell_ref = check.get('cell')
        cell_data = result.get(cell_ref, {})
        cell_score = 0.0
        parts = 0
        matched = 0
        if 'value' in check:
            parts += 1
            if cell_data.get('value') == check['value']:
                matched += 1
        if 'merged' in check:
            parts += 1
            if cell_data.get('merged') == check['merged']:
                matched += 1
        if parts > 0:
            cell_score = matched / parts
        score += per_check * cell_score
    return min(score, 1.0)

def check_xlsx_clevel_count__802aa06dfc0f94fbf79660d5d540798f(result, expected, **options):
    """Check that the C-level count is correct and ranks are split.

    Partial credit:
    - 0.3: E1 has a label like 'C-level Count' or similar
    - 0.3: Rank column (D) is populated for all employees
    - 0.4: E2 contains the correct count of C-level executives
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_count = expected.get('expected_count', 3)
    e1 = result.get('e1_value')
    if e1 is not None and len(str(e1).strip()) > 0:
        score += 0.3
    ranks = result.get('ranks', [])
    non_null_ranks = sum((1 for v in ranks if v is not None))
    if non_null_ranks >= 20:
        score += 0.3
    e2 = result.get('e2_value')
    if e2 is not None:
        try:
            actual_count = int(float(str(e2)))
            if actual_count == expected_count:
                score += 0.4
        except (ValueError, TypeError):
            pass
    return min(score, 1.0)

def check_xlsx_sorted_desc__db00ed878507fde13aefe6b71212bcbd(result, expected, **options):
    """Check if the rows are sorted by citations in descending order."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    rows = result.get('rows', [])
    if len(rows) < 2:
        return 0.0
    expected_names = set(expected.get('expected_names', []))
    actual_names = set((r['name'] for r in rows))
    if expected_names and (not expected_names.issubset(actual_names)):
        return 0.0
    citations = [r['citations'] for r in rows]
    numeric_citations = []
    for c in citations:
        try:
            numeric_citations.append(float(c))
        except (TypeError, ValueError):
            return 0.0
    is_sorted = all((numeric_citations[i] >= numeric_citations[i + 1] for i in range(len(numeric_citations) - 1)))
    return 1.0 if is_sorted else 0.0

def check_cell_value__8afea61a57c3ecda7ab6098e540526c4(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_value):
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_sheet2_merge_text__221a8bb1dbb77d2fdb6346ea5a7576e4(result, expected, **options):
    """Check Sheet2 creation, merge, and text entry with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rules = expected
    if result.get('sheet2_exists'):
        score += 0.34
    if result.get('is_merged'):
        score += 0.33
    expected_text = rules.get('expected_text', 'Demographic Profile')
    actual = result.get('a1_value')
    if actual is not None and str(actual).strip() == expected_text:
        score += 0.33
    return min(score, 1.0)

def check_cell_string_value__04a818774a969e111c0b34c8156dff9c(result, expected, **options):
    """Check if a cell string value matches expected (case-insensitive, stripped)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    if expected_value is None:
        return 0.0
    if str(actual).strip().lower() == str(expected_value).strip().lower():
        return 1.0
    return 0.0

def check_cell_bgcolor__fabb4ab7d265226ac33712a7dcf2e535(result, expected, **options):
    """Check if specified cells have the expected background color.
    Supports partial credit: score = fraction of cells with correct color.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    colors = result.get('colors', {})
    expected_color = expected.get('expected_color', '').upper()
    expected_cells = expected.get('cells', [])
    if not expected_cells or not expected_color:
        return 0.0
    correct = 0
    for cell_ref in expected_cells:
        actual = colors.get(cell_ref)
        if actual and expected_color in actual.upper():
            correct += 1
    return correct / len(expected_cells) if expected_cells else 0.0

def check_cell_value_numeric__85f75e16b65eecbecc9a6c1bce611233(result, expected, **options):
    """Check if a cell's numeric value matches expected within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else None
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.01)
    try:
        actual_f = float(actual)
        expected_f = float(expected_val)
        if abs(actual_f - expected_f) <= tolerance:
            return 1.0
        return 0.0
    except (TypeError, ValueError):
        return 0.0

def check_xlsx_email_cells__93729a4c1a924bf8a2b603d4f55d7635(result, expected, **options):
    """Check if email cells contain expected values. Partial credit per cell."""
    if not result or isinstance(result, str) or result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not expected_values:
        return 0.0
    total = len(expected_values)
    correct = 0
    for (cell_ref, exp_val) in expected_values.items():
        actual = values.get(cell_ref)
        if actual is not None and actual.strip().lower() == str(exp_val).strip().lower():
            correct += 1
    return correct / total if total > 0 else 0.0

def check_cell_value__cb33c763b3ef0ae13364585aa1e6bfa3(result, expected, **options):
    """Check if cell value matches expected value (case-insensitive string comparison)."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None and expected_val is None:
        return 1.0
    if actual is None or expected_val is None:
        return 0.0
    actual_str = str(actual).strip().lower()
    expected_str = str(expected_val).strip().lower()
    if actual_str == expected_str:
        return 1.0
    return 0.0

def check_filtered_sheet__7192ac900176bf2dc148847f36558198(result, expected, **options):
    """Check that a filtered sheet contains the correct movies with proper sorting.

    Partial credit:
    - 0.25: Sheet exists with correct name
    - 0.25: Headers match original format
    - 0.25: Correct number of rows (within tolerance)
    - 0.25: Data is correctly sorted
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_sheet = expected.get('expected_sheet_name', '')
    if result.get('sheet_name', '').strip().lower() == expected_sheet.strip().lower():
        score += 0.25
    expected_headers = expected.get('expected_headers', [])
    actual_headers = result.get('headers', [])
    if len(actual_headers) >= len(expected_headers):
        match = all((str(a).strip().lower() == str(e).strip().lower() for (a, e) in zip(actual_headers, expected_headers)))
        if match:
            score += 0.25
    expected_count = expected.get('expected_row_count', 0)
    actual_count = result.get('row_count', 0)
    if actual_count == expected_count:
        score += 0.25
    elif abs(actual_count - expected_count) <= 2:
        score += 0.125
    sort_key = expected.get('sort_key', '')
    sort_order = expected.get('sort_order', 'desc')
    rows = result.get('rows', [])
    if rows and sort_key:
        values = []
        for r in rows:
            v = r.get(sort_key)
            if v is not None:
                try:
                    values.append(float(str(v)))
                except (ValueError, TypeError):
                    values.append(0)
        if len(values) >= 2:
            if sort_order == 'desc':
                is_sorted = all((values[i] >= values[i + 1] for i in range(len(values) - 1)))
            else:
                is_sorted = all((values[i] <= values[i + 1] for i in range(len(values) - 1)))
            if is_sorted:
                score += 0.25
    return min(score, 1.0)

def check_cell_value__3f903c88c3cac4cb7f6f63b58e6790e4(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
        if abs(actual_num - expected_num) < 0.01:
            return 1.0
    except (TypeError, ValueError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_sheet2_month_avg__d92436572f728d9c9f2cb1a9c19eb76d(result, expected, **options):
    """Check Sheet2 has Month, Total, Average columns with correct values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not result.get('sheet_found'):
        return 0.0
    score = 0.0
    headers = [h.lower() for h in result.get('headers', [])]
    expected_data = expected.get('expected_data', {})
    num_months = len(expected_data)
    required_headers = {'month', 'total', 'average'}
    if required_headers.issubset(set(headers)):
        score += 0.2
    total_correct = 0
    for (month, vals) in expected_data.items():
        actual_row = result.get('data', {}).get(month, {})
        for (key, val) in actual_row.items():
            if key.lower() == 'total':
                if val is not None and abs(float(val) - vals['total']) < 1:
                    total_correct += 1
                break
    if num_months > 0:
        score += 0.4 * (total_correct / num_months)
    avg_correct = 0
    for (month, vals) in expected_data.items():
        actual_row = result.get('data', {}).get(month, {})
        for (key, val) in actual_row.items():
            if key.lower() == 'average':
                if val is not None and abs(float(val) - vals['average']) < 2:
                    avg_correct += 1
                break
    if num_months > 0:
        score += 0.4 * (avg_correct / num_months)
    return min(score, 1.0)

def check_cells_b3_b6__f927596325f5c4e25b12a463bfafbace(result, expected, **options):
    """Check if cells B3-B6 all have the expected value, with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_val = expected.get('expected_value')
    cells = expected.get('cells', ['B3', 'B4', 'B5', 'B6'])
    if not expected_val or not values:
        return 0.0
    score = 0.0
    per_cell = 1.0 / len(cells)
    for cell in cells:
        actual = values.get(cell)
        if actual is not None and str(actual).strip().lower() == str(expected_val).strip().lower():
            score += per_cell
    return min(score, 1.0)

def check_xls_cell_value__2674782c2441965e8475070c5235e39b(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if float(actual) == float(expected_val):
            return 1.0
    except (ValueError, TypeError):
        pass
    if str(actual).strip() == str(expected_val).strip():
        return 1.0
    return 0.0

def check_header_cells__25536421a644aed93ef4482447c70e86(result, expected, **options):
    """Check if E1 and F1 headers were added correctly. Partial credit per header."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_e1 = expected.get('expected_e1', '').lower().strip()
    expected_f1 = expected.get('expected_f1', '').lower().strip()
    actual_e1 = (result.get('e1') or '').lower().strip()
    actual_f1 = (result.get('f1') or '').lower().strip()
    if actual_e1 and expected_e1 in actual_e1:
        score += 0.5
    if actual_f1 and expected_f1 in actual_f1:
        score += 0.5
    return min(score, 1.0)

def check_pptx_table_cell__1546f90c5e4d91a5767b01e5f0a56119(result, expected, **options):
    """Check if a table cell text matches expected text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_text = result.get('cell_text', '').strip()
    expected_text = expected.get('expected_text', '').strip()
    if actual_text == expected_text:
        return 1.0
    if actual_text.lower() == expected_text.lower():
        return 0.8
    return 0.0

def check_cell_numeric__13d011eac188a961b7fc2a44b3a2069f(result, expected, **options):
    """Check if a cell's numeric value matches expected."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        actual_num = float(actual)
        expected_num = float(expected_value)
        if abs(actual_num - expected_num) < 0.001:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_cell_numeric_close__7aae35559065c2f6c82a9e66b7d55095(result, expected, **options):
    """Check if a cell value is numerically close to expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    value = result.get('value')
    if value is None:
        return 0.0
    try:
        actual = float(value)
    except (TypeError, ValueError):
        return 0.0
    expected_val = expected.get('expected_value')
    if expected_val is None:
        return 0.0
    tolerance = expected.get('tolerance', 1.0)
    if abs(actual - float(expected_val)) <= tolerance:
        return 1.0
    return 0.0

def check_cell_value__cb5b990d93f87d0f04a6150ee1d23a05(result, expected, **options):
    """Check if a cell value matches expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None and expected_val is None:
        return 1.0
    if actual is not None and expected_val is not None:
        if str(actual).strip().lower() == str(expected_val).strip().lower():
            return 1.0
    return 0.0

def check_multi_cell_values__7424fba2d2a12056c33da3db49551203(result, expected, **options):
    """Check multiple cell values with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    values = result.get('values', {})
    if not values:
        return 0.0
    checks = expected.get('checks', [])
    if not checks:
        return 0.0
    total_weight = sum((c.get('weight', 1.0) for c in checks))
    if total_weight == 0:
        return 0.0
    score = 0.0
    for check in checks:
        cell = check.get('cell')
        expected_val = check.get('expected')
        weight = check.get('weight', 1.0)
        tolerance = check.get('tolerance', 0.01)
        actual = values.get(cell)
        if actual is None and expected_val is None:
            score += weight
            continue
        if actual is None or expected_val is None:
            continue
        if isinstance(expected_val, str):
            if str(actual).strip().lower() == expected_val.strip().lower():
                score += weight
            continue
        try:
            actual_num = float(actual)
            expected_num = float(expected_val)
            if abs(actual_num - expected_num) <= tolerance:
                score += weight
        except (TypeError, ValueError):
            pass
    return min(score / total_weight, 1.0)

def check_cell_bgcolor__0e55cf8187f0e07889b1109cdd0b266f(result, expected, **options):
    """Check if specified cells have the expected background color.
    Supports partial credit: score = fraction of cells with correct color.
    """
    if isinstance(result, str) or result.get('error'):
        return 0.0
    colors = result.get('colors', {})
    expected_color = expected.get('expected_color', '').upper()
    expected_cells = expected.get('cells', [])
    if not expected_cells or not expected_color:
        return 0.0
    correct = 0
    for cell_ref in expected_cells:
        actual = colors.get(cell_ref)
        if actual and expected_color in actual.upper():
            correct += 1
    return correct / len(expected_cells) if expected_cells else 0.0

def check_cell_value__990c97f8e62011f017525e7c1f376ba9(result, expected, **options):
    """Check if a cell value matches the expected value."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    actual_str = str(actual).strip()
    expected_str = str(expected_value).strip()
    if actual_str == expected_str:
        return 1.0
    return 0.0

def check_xlsx_qty_sort__d422460342c18f4ab4e35ab11ffce7d8(result, expected, **options):
    """Check if spreadsheet data is sorted by Quantity descending."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_sorted_descending'):
        score += 0.7
    expected_rows = expected.get('expected_row_count', 35)
    if result.get('row_count') == expected_rows:
        score += 0.15
    expected_first = expected.get('expected_first_quantity', 124)
    if result.get('first_quantity') == expected_first:
        score += 0.15
    return min(score, 1.0)

def check_cell_value__949a54a656a59ff2d23607cb949dacfc(result, expected, **options):
    """Check if a single cell value matches expected value with numeric tolerance."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    try:
        if abs(float(actual) - float(expected_val)) < 0.5:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_cell_value_numeric__28ae1ff6307904e177d5ad2de88ea142(result, expected, **options):
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.5)
    try:
        actual_num = float(actual)
        expected_num = float(expected_val)
    except (TypeError, ValueError):
        return 0.0
    if abs(actual_num - expected_num) <= tolerance:
        return 1.0
    return 0.0

def check_sheet_names__a164952ae6b41142a59183faf6bedadf_qw35sft2_d3debb98(result, expected, **options):
    """Check that the sheet names match exactly in the correct order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('sheet_names', [])
    expected_names = expected.get('sheet_names', [])
    if actual == expected_names:
        return 1.0
    return 0.0

def check_merged_cell_value__48a88b5bb3b362cc9c78cc671cedfc7c_qw35sft2_7a9956ec(result, expected, **options):
    """Check merged cell and its text value. Partial credit: 0.5 for merge, 0.5 for value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('is_merged'):
        score += 0.5
    actual_value = result.get('value')
    expected_value = expected.get('expected_value', '')
    if actual_value is not None and str(actual_value).strip() == str(expected_value).strip():
        score += 0.5
    return score

def check_cell_numeric_value__0762adbae4ad235c185f47096cf64c91_qw35sft2_0a2abca9(result, expected, **options):
    """Check that a cell contains the expected integer count (accepts both int and string forms)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_count = expected.get('expected_count')
    if expected_count is None:
        return 0.0
    try:
        return 1.0 if int(actual) == int(expected_count) else 0.0
    except (ValueError, TypeError):
        return 0.0

def check_sheet2_state__68824c07a6606ac16add37bf4765401b_qw35sft2_5cf70407(result, expected, **options):
    """Partial credit: 0.5 for Sheet2 existence, 0.5 for correct A1 text."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('sheet_exists', False):
        score += 0.5
    expected_text = expected.get('a1_value', 'Demographic Profile')
    actual = result.get('a1_value')
    if actual is not None and str(actual).strip() == str(expected_text).strip():
        score += 0.5
    return score

def check_seqno_and_sheet_name__83227706efd58da00bb8213cee393c7a_qw35sft2_523e02bd(result, expected, **options):
    """Check Seq No. column (B2:B29) is correct and the sheet is renamed as expected."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    seq_nos = result.get('seq_nos', [])
    expected_seq = [f'No. {i}' for i in range(1, 29)]
    if seq_nos == expected_seq:
        score += 0.5
    expected_name = expected.get('expected_sheet_name', 'Sales Data')
    if result.get('sheet_name') == expected_name:
        score += 0.5
    return score

def check_cell_value__cfba55e4566ed0373f9b47631b661201_qw35sft2_79ef5ec5(result, expected, **options):
    """Check if a cell contains the expected string value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None or expected_value is None:
        return 0.0
    if isinstance(expected_value, str):
        return 1.0 if str(actual).strip() == expected_value.strip() else 0.0
    return 1.0 if actual == expected_value else 0.0

def check_len_formula__29aa3f16a2f40add619bb0073f57ed2e_qw35sft2_55c5c202(result, expected, **options):
    """Check C2 has correct cleaned title (0.5) and D2 has correct LEN value (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    c2_actual = result.get('c2_value')
    c2_expected = expected.get('c2_expected')
    if c2_actual is not None and c2_expected is not None:
        if str(c2_actual).strip() == str(c2_expected).strip():
            score += 0.5
    d2_actual = result.get('d2_value')
    d2_expected = expected.get('d2_expected')
    if d2_actual is not None and d2_expected is not None:
        try:
            if int(d2_actual) == int(d2_expected):
                score += 0.5
        except (ValueError, TypeError):
            pass
    return score

def check_cell_value_equals__d4b207925bd2d0c3ababe1e319fa8b1f_qw35sft2_5621d4f6(result, expected, **options):
    """Return 1.0 if the cell value matches expected_value (numeric or string)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else result
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    if isinstance(expected_val, (int, float)) and isinstance(actual, (int, float)):
        return 1.0 if abs(float(actual) - float(expected_val)) < 0.01 else 0.0
    if isinstance(expected_val, str) and isinstance(actual, str):
        return 1.0 if actual.strip().lower() == expected_val.strip().lower() else 0.0
    return 1.0 if actual == expected_val else 0.0

def check_calc_earned__73e90374b4fd7034756570a58d380c35_qw35sft2_61998100(result, expected, **options):
    """Check that E3 contains a value close to expected total earned (191.666...)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual = result.get('e3_value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value', 191.6667)
    tolerance = expected.get('tolerance', 0.01)
    return 1.0 if abs(float(actual) - float(expected_value)) <= tolerance else 0.0

def check_na_count_in_cell__1ed2b6f96cb9e3177c42c0765213bb8c_qw35sft2_65832f34(result, expected, **options):
    """Check that cell A1 contains the expected NA row count (13)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual_val = result.get('value')
    expected_val = expected.get('expected_value', 13)
    if actual_val is None:
        return 0.0
    try:
        return 1.0 if int(actual_val) == int(expected_val) else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_multi_cell_grade__3b5b28cd38bc58813bd97e9b2f5d4ad4_qw35sft2_e971eac3(result, expected, **options):
    """Check multiple cells for correct grade values with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    cell_grades = expected.get('cell_grades', {})
    if not cell_grades:
        return 0.0
    score = 0.0
    weight = 1.0 / len(cell_grades)
    for cell, expected_val in cell_grades.items():
        actual = result.get(cell)
        if actual is not None and str(actual).strip() == str(expected_val).strip():
            score += weight
    return min(score, 1.0)

def check_avg_age_cell__c483fe0106f1e0b2430ad19cd97c77c0_qw35sft2_18c5427e(result, expected, **options):
    """Check that cell E3 contains a numeric value close to the expected average age."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    val = result.get('E3')
    if val is None:
        return 0.0
    try:
        return 1.0 if abs(float(val) - float(expected.get('expected_value', 33.0))) < 0.5 else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_calc_sheet1_net_income__6142d421fefd6b784b8ad81070a58350_qw35sft2_9fa35d28(result, expected, **options):
    """Check Sheet1 column C has 'Net Income' header and correct per-row values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    exp_header = expected.get('header', 'Net Income')
    exp_values = expected.get('values', [29101, 40412, 23036, 20782, 27289, 7980, 26615, 31546, 14631, 1969, 7348, 34004, 19438, 25326, 15863, 39308, 27916, 30460, 33076])
    actual_header = result.get('header')
    if isinstance(actual_header, str) and actual_header.strip() == exp_header.strip():
        score += 0.2
    actual_values = result.get('values', [])
    per_row_weight = 0.8 / len(exp_values)
    for i, exp_v in enumerate(exp_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if isinstance(actual, float) and actual == int(actual):
                actual = int(actual)
            if actual == exp_v:
                score += per_row_weight
    return round(min(score, 1.0), 4)

def check_sheet2_sorted_revenue__ff6f4e09739d1959467ff60470ffe2bd_qw35sft2_dae27611(result, expected, **options):
    """Check Sheet2 Revenue column exists and is sorted in ascending order."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('header') == expected.get('header'):
        score += 0.2
    if result.get('row_count') == expected.get('row_count'):
        score += 0.2
    if result.get('is_sorted_asc'):
        score += 0.4
    if result.get('first_value') == expected.get('min_value'):
        score += 0.1
    if result.get('last_value') == expected.get('max_value'):
        score += 0.1
    return min(score, 1.0)

def check_cell_string__9ab5f66fb165e5668bd5b38cac0c73c2_qw35sft2_0f40b827(result, expected, **options):
    """Check if the cell value matches the expected string."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value', '').strip() if isinstance(result, dict) else str(result).strip()
    expected_value = expected.get('expected_value', '').strip()
    return 1.0 if actual == expected_value else 0.0

def check_cells_red__03a99dd86c6d30983d2467bc1177489c_qw35sft2_ef4d5f4f(result, expected, **options):
    """Check that all target cells have the expected red background color.
    Returns partial credit based on fraction of target cells correctly colored."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    target_cells = expected.get('target_cells', [])
    expected_color = expected.get('expected_color', 'FFFF0000').upper()
    if not target_cells:
        return 0.0
    correct = sum((1 for cell in target_cells if result.get(cell, '').upper() == expected_color))
    return correct / len(target_cells)

def check_xlsx_transposed_cells__f7593e4503a45a853c1ae96cb08aaf92_qw35sft2_b1c16c4d(result, expected, **options):
    """Check that key cells in the transposed table B8:E12 match expected values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = [('B8', expected.get('B8')), ('E8', expected.get('E8')), ('B9', expected.get('B9')), ('E9', expected.get('E9')), ('B12', expected.get('B12')), ('E12', expected.get('E12'))]
    correct = 0
    for cell, exp_val in checks:
        actual = result.get(cell)
        if exp_val is None:
            continue
        try:
            if isinstance(exp_val, (int, float)) and actual is not None:
                if abs(float(actual) - float(exp_val)) < 0.01:
                    correct += 1
                    continue
        except (TypeError, ValueError):
            pass
        if actual == exp_val:
            correct += 1
    return correct / len(checks)

def check_calc_sheet_renamed__ecf996871673b767fb6d252021f29e29_qw35sft2_cf744cd5(result, expected, **options):
    """Check that Sheet1 was renamed to the new name. Full credit only if old name is gone."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    new_name = expected.get('new_sheet_name', 'SalesData')
    old_name = expected.get('old_sheet_name', 'Sheet1')
    sheet_names = result.get('sheet_names', [])
    has_new = new_name in sheet_names
    lacks_old = old_name not in sheet_names
    if has_new and lacks_old:
        return 1.0
    elif has_new:
        return 0.5
    return 0.0

def check_xlsx_zone2_row_totals__6ef64aee01b9e4fef352e390ba82e0ce_qw35sft2_d34b1ef4(result, expected, **options):
    """Check F10, F11, F12 match expected Zone 2 product totals with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    checks = []
    for cell_key, exp_key in [('F10', 'F10_expected'), ('F11', 'F11_expected'), ('F12', 'F12_expected')]:
        actual = result.get(cell_key)
        exp = expected.get(exp_key)
        try:
            ok = actual is not None and abs(float(actual) - float(exp)) < 0.5
        except (TypeError, ValueError):
            ok = False
        checks.append(ok)
    return sum(checks) / len(checks)

def check_sheet2_three_columns__74fb6bc2ea3e8707d0a1b3dd5202c02d_qw35sft2_4ffd667d(result, expected, **options):
    """Check Sheet2 has Month/Total/Average columns with correct values (partial credit)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    h_a = str(result.get('header_a1', '') or '').strip().lower()
    h_b = str(result.get('header_b1', '') or '').strip().lower()
    h_c = str(result.get('header_c1', '') or '').strip().lower()
    if h_a == 'month' and h_b == 'total':
        score += 0.15
    if h_c in ('average', 'avg'):
        score += 0.1
    month_totals = expected.get('month_totals', {})
    months = result.get('months', [])
    totals = result.get('totals', [])
    month_map = {str(m): t for m, t in zip(months, totals) if m}
    correct_totals = sum((1 for m, ev in month_totals.items() if month_map.get(m) is not None and abs(float(month_map[m]) - float(ev)) <= 1))
    if correct_totals == 6:
        score += 0.35
    elif correct_totals >= 3:
        score += 0.18
    averages_expected = expected.get('averages', {})
    averages_got = result.get('averages', [])
    avg_map = {str(m): a for m, a in zip(months, averages_got) if m}
    correct_avgs = sum((1 for m, ea in averages_expected.items() if avg_map.get(m) is not None and abs(float(avg_map[m]) - float(ea)) <= 1))
    if correct_avgs == 6:
        score += 0.4
    elif correct_avgs >= 3:
        score += 0.2
    return min(score, 1.0)

def check_ramp_accel_cells__f874da4b0492e2e023290240d31f1c8d_qw35sft2_083dc380(result, expected, **options):
    """Check that columns B and D are filled with correct acceleration values.

    Partial credit: 0.25 each for B10, B30, D10, D30 within tolerance 0.015.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tol = 0.015
    score = 0.0
    _near = lambda actual, exp: actual is not None and abs(float(actual) - float(exp)) <= tol
    if _near(result.get('b10'), expected.get('b10_expected')):
        score += 0.25
    if _near(result.get('b30'), expected.get('b30_expected')):
        score += 0.25
    if _near(result.get('d10'), expected.get('d10_expected')):
        score += 0.25
    if _near(result.get('d30'), expected.get('d30_expected')):
        score += 0.25
    return round(score, 4)

def check_sheet2_pct_format__6c9c981a9ca6d475ebfb8637fc973b0b_qw35sft2_c08948ee(result, expected, **options):
    """
    Check Sheet2 has correct % change values AND cells B2:D6 are formatted as percentage.
    Partial credit: 0.5 for correct values, 0.5 for percentage format on B2:D6.
    """
    if not result or result.get('error'):
        return 0.0
    score = 0.0
    tol = 0.001
    headers = result.get('headers', [])
    expected_headers = ['Year', 'CA changes', 'FA changes', 'OA changes']
    headers_ok = all((str(h).strip() == e for h, e in zip(headers, expected_headers))) if headers else False
    data_rows = result.get('data_rows', [])
    expected_years = [2015, 2016, 2017, 2018, 2019]
    expected_ca = [0.10149072, 0.07217629, 0.1342019, 0.06462741, 0.06554579]
    expected_fa = [-0.0496044, -0.05556969, -0.05928012, -0.06676384, -0.07455235]
    expected_oa = [-0.01675978, 0.05852273, 0.07648953, 0.00473697, 0.01439206]
    data_ok = False
    if len(data_rows) >= 5 and headers_ok:
        rows_correct = 0
        for i, row in enumerate(data_rows[:5]):
            if row[0] is None:
                break
            try:
                year_ok = int(row[0]) == expected_years[i]
                ca_ok = row[1] is not None and abs(float(row[1]) - expected_ca[i]) < tol
                fa_ok = row[2] is not None and abs(float(row[2]) - expected_fa[i]) < tol
                oa_ok = row[3] is not None and abs(float(row[3]) - expected_oa[i]) < tol
                if year_ok and ca_ok and fa_ok and oa_ok:
                    rows_correct += 1
            except (TypeError, ValueError):
                pass
        data_ok = rows_correct >= 4
    if headers_ok and data_ok:
        score += 0.5
    formats = result.get('formats', [])
    if len(formats) >= 5:
        pct_cells = 0
        total_cells = 0
        for row_fmts in formats[:5]:
            for fmt in row_fmts:
                total_cells += 1
                if fmt is not None and '%' in str(fmt).lower():
                    pct_cells += 1
        if total_cells > 0 and pct_cells / total_cells >= 0.8:
            score += 0.5
    return min(score, 1.0)

def check_xlsx_padded_max__7fad2dd93f80fb906237dbbf8bdc9fb2_qw35sft2_527a04cd(result, expected, **options):
    """Check D column zero-padding (0.5) and E1 contains max Old ID value (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rows = result.get('rows', [])
    if rows:
        correct = sum((1 for r in rows if r.get('d') == str(r['c']).zfill(7)))
        if correct == len(rows):
            score += 0.5
    max_expected = expected.get('max_expected', 21540)
    e1 = result.get('e1')
    try:
        if e1 is not None and int(e1) == max_expected:
            score += 0.5
    except (ValueError, TypeError):
        pass
    return score

def check_xlsx_spent_and_date_format__2b619f85a4e6da743c7b9581de6419ad_qw35sft2_99cbf667(result, expected, **options):
    """Check C2:C8 have 2dp format and B2:B8 have YYYY-MM-DD date format."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    c_formats = result.get('c_formats', [])
    has_two_dp = bool(c_formats) and all((fmt is not None and ('0.00' in str(fmt) or '##0.00' in str(fmt)) for fmt in c_formats))
    if has_two_dp:
        score += 0.5
    b_formats = result.get('b_formats', [])
    has_yyyy_mmdd = bool(b_formats) and all((fmt is not None and 'yyyy' in str(fmt).lower() and ('mm' in str(fmt).lower()) and ('dd' in str(fmt).lower()) for fmt in b_formats))
    if has_yyyy_mmdd:
        score += 0.5
    return score

def check_calc_pivot_and_sort__6bb297a00c63757f69bb3bc219190d5b_qw35sft2_5e630c32(result, expected, **options):
    """
    Partial credit:
      0.5 - Sheet2 exists with correct pivot table (invoice 10505 count == 5)
      0.5 - Sheet1.G2 == 750 (max Sales value, first after descending sort)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    sheet2_exists = result.get('sheet2_exists', False)
    count_10505 = result.get('invoice_count_10505')
    if sheet2_exists:
        try:
            if int(count_10505) == 5:
                score += 0.5
            else:
                score += 0.25
        except (TypeError, ValueError):
            score += 0.25
    expected_first_sale = expected.get('first_sale', 750)
    sheet1_g2 = result.get('sheet1_g2')
    try:
        if abs(float(sheet1_g2) - float(expected_first_sale)) < 0.01:
            score += 0.5
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_cell_string__13e2a9cff490c6125c903160099b9b7a_qw35sft2_d6b97dfc(result, expected, **options):
    """Partial-credit evaluation for the three-part pivot-table task.

    Scoring breakdown:
      0.5 — Sheet2 exists, has >=8 non-empty cells, contains the best-selling
             product name (product-revenue pivot confirmed), AND contains at least
             one known sales channel label (channel-revenue pivot confirmed).
      0.5 — Sheet1!I1 equals the expected best-selling product name.
    Total: 1.0
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    expected_val = str(expected.get('value', '')).strip()
    sheet2_exists = result.get('sheet2_exists', False)
    sheet2_nonempty = result.get('sheet2_nonempty', 0)
    sheet2_str_values = result.get('sheet2_str_values', [])
    known_channels = {'E-mail Coupon', 'In Store Sales', 'Web Site Sales'}
    channel_pivot_present = any((ch in sheet2_str_values for ch in known_channels))
    pivot_tables_present = sheet2_exists and sheet2_nonempty >= 8 and (expected_val in sheet2_str_values) and channel_pivot_present
    if pivot_tables_present:
        score += 0.5
    actual = result.get('value')
    if actual is not None and str(actual).strip() == expected_val:
        score += 0.5
    return min(score, 1.0)

def check_month_pivot_sheet2__360bd21f4036eac9068c6b83aba7e2e2_qw35sft2_a4406367(result, expected, **options):
    """Check that Sheet2 contains expected monthly revenue totals and month name labels."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    if not result.get('sheet2_exists'):
        return 0.0
    tolerance = expected.get('tolerance', 1.0)
    numeric_values = result.get('numeric_values', [])
    string_values = result.get('string_values', [])
    expected_months = expected.get('expected_months', ['May', 'Jun', 'Jul', 'Aug', 'Sep'])
    months_found = sum((1 for m in expected_months if any((m.lower() in s.lower() for s in string_values))))
    month_score = months_found / len(expected_months) if expected_months else 0.0
    expected_revenues = [float(r) for r in expected.get('expected_revenues', [])]
    if expected_revenues:
        rev_found = sum((1 for exp_rev in expected_revenues if any((abs(v - exp_rev) <= tolerance for v in numeric_values))))
        rev_score = rev_found / len(expected_revenues)
    else:
        rev_score = 0.0
    return 0.4 * month_score + 0.6 * rev_score

def check_freeze_and_sheet__c058d8a3342bf3e7968eb49c80bec94d_qw35sft2_c4c06d2f(result, expected, **options):
    """Check freeze_panes == 'C2' (0.5) and sheet name matches expected (0.5)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_freeze = expected.get('freeze_panes', 'C2')
    if result.get('freeze_panes') == expected_freeze:
        score += 0.5
    expected_sheet = expected.get('sheet_name', 'Monthly Data')
    if result.get('sheet_name') == expected_sheet:
        score += 0.5
    return score

def check_cell_numeric__e7e54b6e523b2a67e0ff77298aeb57dc_qw35sft2_c16bee5a(result, expected, **options):
    """Check if cell value matches expected integer (tolerates int/float difference)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value')
    try:
        if int(float(actual)) == int(expected_value):
            return 1.0
    except (TypeError, ValueError):
        pass
    return 1.0 if actual == expected_value else 0.0

def check_cell_numeric__57a509452fa7aaef388db657e08eeb87_qw35sft2_462d120f(result, expected, **options):
    """Check if a cell value is within tolerance of the expected numeric value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    tolerance = expected.get('tolerance', 0.1)
    if actual is None or expected_value is None:
        return 0.0
    try:
        return 1.0 if abs(float(actual) - float(expected_value)) <= tolerance else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_cell_value_equals__68e5aa381affaca2d81c1d380acacc93_qw35sft2_c346be8b(result, expected, **options):
    """Return 1.0 if the cell value matches expected_value (numeric or string)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else result
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    if isinstance(expected_val, (int, float)) and isinstance(actual, (int, float)):
        return 1.0 if abs(float(actual) - float(expected_val)) < 0.01 else 0.0
    if isinstance(expected_val, str) and isinstance(actual, str):
        return 1.0 if actual.strip().lower() == expected_val.strip().lower() else 0.0
    return 1.0 if actual == expected_val else 0.0

def check_cell_value__009b0a2429b483dc78b74a0168744c1a_qw35sft2_f16b38bd(result, expected, **options):
    """Check that a cell value matches expected_value (string comparison, strip whitespace)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value', '')
    return 1.0 if str(actual).strip() == str(expected_value).strip() else 0.0

def check_multi_cell_text__142de785ea8e1e7c66581a5013a34ff3_qw35sft2_99cabfcf(result, expected, **options):
    """Check multiple cell values with partial credit. Each cell contributes equally."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    values = result.get('values', {})
    expected_values = expected.get('expected_values', {})
    if not expected_values:
        return 0.0
    n = len(expected_values)
    score = 0.0
    for cell_ref, exp_val in expected_values.items():
        actual = values.get(cell_ref)
        if actual is not None and str(actual).strip() == str(exp_val).strip():
            score += 1.0 / n
    return round(min(score, 1.0), 4)

def check_sheet_names_partial__051a1a86398cf9c50b041fec4af58dd3_qw35sft2_2917bf35(result, expected, **options):
    """Partial-credit check for 4-sheet rename+copy+delete scenario.

    Scoring:
    - 0.75: first three sheets are ['LARS Resources', 'LARS Resources (Backup)', 'LARS Resources (Offline)']
    - 0.25: fourth sheet is 'LARS Archive'
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('sheet_names', [])
    expected_names = expected.get('sheet_names', [])
    if not expected_names:
        return 0.0
    score = 0.0
    if len(actual) >= 3 and actual[:3] == expected_names[:3]:
        score += 0.75
    if len(actual) >= 4 and actual[3] == expected_names[3]:
        score += 0.25
    return min(score, 1.0)

def check_cell_grade__b0dcd11a1bdcc1e13016c1ab56479e5f_qw35sft2_0d0437cb(result, expected, **options):
    """Check that a cell contains the expected grade string."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else result
    expected_val = expected.get('expected_value')
    if actual is None or expected_val is None:
        return 0.0
    return 1.0 if str(actual).strip() == str(expected_val).strip() else 0.0

def check_cell_string__76b216d4ab5877594f0cff145afca3ba_qw35sft2_75ae6314(result, expected, **options):
    """Check if the cell value matches the expected string."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value', '').strip() if isinstance(result, dict) else str(result).strip()
    expected_value = expected.get('expected_value', '').strip()
    return 1.0 if actual == expected_value else 0.0

def check_calc_summary_sheet__878e408cd8149dcd33226991d9a64c87_qw35sft2_b8e83f56(result, expected, **options):
    """Check that 'Summary' sheet exists with correct headers and sum values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('sheet_exists'):
        score += 0.25
    if isinstance(result.get('A1'), str) and result['A1'].strip() == expected.get('header_revenue', 'Total Revenue'):
        score += 0.125
    if isinstance(result.get('B1'), str) and result['B1'].strip() == expected.get('header_expenses', 'Total Expenses'):
        score += 0.125
    a2 = result.get('A2')
    if isinstance(a2, float) and a2 == int(a2):
        a2 = int(a2)
    if a2 == expected.get('sum_revenue', 867786):
        score += 0.25
    b2 = result.get('B2')
    if isinstance(b2, float) and b2 == int(b2):
        b2 = int(b2)
    if b2 == expected.get('sum_expenses', 411686):
        score += 0.25
    return round(min(score, 1.0), 4)

def check_sheet1_net_income_sheet2_revenue__5d938927d392c1d220c2599bb703f167_qw35sft2_63e26994(result, expected, **options):
    """Check Sheet2 has Revenue column AND Sheet1 has a Net Income column C."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('sheet2_header') == 'Revenue' and result.get('sheet2_row_count') == 19:
        score += 0.4
    c_header = result.get('sheet1_c_header', '')
    if c_header and 'net' in str(c_header).lower() and ('income' in str(c_header).lower()):
        score += 0.2
    if result.get('sheet1_c_row_count') == 19:
        score += 0.2
    expected_sum = expected.get('net_income_sum')
    actual_sum = result.get('sheet1_c_sum')
    if actual_sum is not None and expected_sum is not None:
        try:
            if abs(float(actual_sum) - float(expected_sum)) < 1:
                score += 0.2
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_calc_total_sales__dc5d11c569bf87971ba13d9e911e3365_qw35sft2_1d6bedbb(result, expected, **options):
    """Check B12 contains the correct total of all Sales values (901745)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_value = expected.get('expected_value')
    actual_value = result.get('value')
    if actual_value is None:
        return 0.0
    try:
        tolerance = expected.get('tolerance', 0.5)
        if abs(float(actual_value) - float(expected_value)) <= tolerance:
            return 1.0
    except (TypeError, ValueError):
        pass
    return 0.0

def check_calc_earned_and_tax__2412d7535e20b0348d2c25accb946746_qw35sft2_41466692(result, expected, **options):
    """Partial credit: 0.5 for correct E3 (total earned ~191.67), 0.5 for correct E4 (20% tax ~38.33)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tolerance = expected.get('tolerance', 0.01)
    e3 = result.get('e3_value')
    expected_e3 = expected.get('expected_e3', 191.6667)
    if e3 is not None and abs(float(e3) - float(expected_e3)) <= tolerance:
        score += 0.5
    e4 = result.get('e4_value')
    expected_e4 = expected.get('expected_e4', 38.3333)
    if e4 is not None and abs(float(e4) - float(expected_e4)) <= tolerance:
        score += 0.5
    return score

def check_cells_red__1a5d60ee39c11d25031d40117b38947f_qw35sft2_b467b721(result, expected, **options):
    """Check that all target cells have the expected red background color.
    Returns partial credit based on fraction of target cells correctly colored."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    target_cells = expected.get('target_cells', [])
    expected_color = expected.get('expected_color', 'FFFF0000').upper()
    if not target_cells:
        return 0.0
    correct = sum((1 for cell in target_cells if result.get(cell, '').upper() == expected_color))
    return correct / len(target_cells)

def check_sheet_renamed__0c20409273ae865247912386f480b246_qw35sft2_73a65b87(result, expected, **options):
    """Check that the sheet was renamed to the expected name."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    expected_name = expected.get('expected_name', 'Data')
    sheet_names = result.get('sheet_names', [])
    if expected_name in sheet_names:
        return 1.0
    return 0.0

def check_xlsx_transposed_sorted__f581e5220b9d6fb8293dc4e0f669a755_qw35sft2_1d570f56(result, expected, **options):
    """Check transposed table is sorted by Marks descending: Olivia(83), Clint(36), Parah(35), Elijah(30)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('B8') == expected.get('B8', 'Student'):
        score += 0.2
    rows = [('B9', 'E9', expected.get('B9', 'Olivia'), expected.get('E9', 83)), ('B10', 'E10', expected.get('B10', 'Clint'), expected.get('E10', 36)), ('B11', 'E11', expected.get('B11', 'Parah'), expected.get('E11', 35)), ('B12', 'E12', expected.get('B12', 'Elijah'), expected.get('E12', 30))]
    for name_cell, marks_cell, exp_name, exp_marks in rows:
        row_correct = 0
        if result.get(name_cell) == exp_name:
            row_correct += 1
        actual_marks = result.get(marks_cell)
        if actual_marks is not None:
            try:
                if abs(float(actual_marks) - float(exp_marks)) < 0.01:
                    row_correct += 1
            except (TypeError, ValueError):
                pass
        score += 0.2 * (row_correct / 2)
    return min(score, 1.0)

def check_xlsx_zone1_row_totals__3aa70d99f23178c4e0eb339ad1e35c2e_qw35sft2_9d875225(result, expected, **options):
    """Check F3, F4, F5 match expected Zone 1 product totals with partial credit."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    for cell, key in [('F3', 'F3_expected'), ('F4', 'F4_expected'), ('F5', 'F5_expected')]:
        actual = result.get(cell)
        exp = expected.get(key)
        if actual is not None and exp is not None:
            try:
                if abs(float(actual) - float(exp)) < 0.5:
                    score += 1.0
            except (TypeError, ValueError):
                pass
    return score / 3.0

def check_sheet2_with_grand_total__0378b78b340823f75ef2df5ebaa9b121_qw35sft2_947a6f38(result, expected, **options):
    """Check Sheet2 has Month/Total headers, correct monthly totals, and a Grand Total row."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    h_a = str(result.get('header_a1', '') or '').strip()
    h_b = str(result.get('header_b1', '') or '').strip()
    if h_a.lower() == 'month' and h_b.lower() == 'total':
        score += 0.2
    month_totals = expected.get('month_totals', {})
    months_in_result = result.get('months', [])
    totals_in_result = result.get('totals', [])
    month_map = {}
    for m, t in zip(months_in_result, totals_in_result):
        if m:
            month_map[str(m).strip()] = t
    correct_months = 0
    for month, expected_total in month_totals.items():
        actual = month_map.get(month)
        if actual is not None and abs(float(actual) - float(expected_total)) <= 1:
            correct_months += 1
    if correct_months == 6:
        score += 0.4
    elif correct_months >= 3:
        score += 0.2
    label = str(result.get('row8_label', '') or '').strip().lower()
    grand_total_expected = float(expected.get('grand_total', 26159))
    grand_val = result.get('row8_value')
    label_ok = 'grand' in label or 'total' in label
    val_ok = grand_val is not None and abs(float(grand_val) - grand_total_expected) <= 1
    if label_ok and val_ok:
        score += 0.4
    elif val_ok:
        score += 0.2
    return min(score, 1.0)

def check_sheet2_sorted__2fc0837f7c43e49c220e69070cca644b_qw35sft2_628d001d(result, expected, **options):
    """
    Check Sheet2 has correct headers, correct change values, and data sorted
    by CA changes in descending order.
    Partial credit: 0.5 for correct values, 0.5 for correct sort order.
    """
    if not result or result.get('error'):
        return 0.0
    score = 0.0
    tol = 0.001
    headers = result.get('headers', [])
    expected_headers = ['Year', 'CA changes', 'FA changes', 'OA changes']
    headers_ok = all((str(h).strip() == e for h, e in zip(headers, expected_headers))) if headers else False
    data_rows = result.get('data_rows', [])
    expected_sorted_years = [2017, 2015, 2016, 2019, 2018]
    year_to_vals = {2015: (0.10149072, -0.0496044, -0.01675978), 2016: (0.07217629, -0.05556969, 0.05852273), 2017: (0.1342019, -0.05928012, 0.07648953), 2018: (0.06462741, -0.06676384, 0.00473697), 2019: (0.06554579, -0.07455235, 0.01439206)}
    if len(data_rows) >= 5 and headers_ok:
        found_years = set()
        all_values_ok = True
        for row in data_rows[:5]:
            if row[0] is None:
                all_values_ok = False
                break
            try:
                yr = int(row[0])
                if yr not in year_to_vals:
                    all_values_ok = False
                    break
                exp_ca, exp_fa, exp_oa = year_to_vals[yr]
                ca_ok = row[1] is not None and abs(float(row[1]) - exp_ca) < tol
                fa_ok = row[2] is not None and abs(float(row[2]) - exp_fa) < tol
                oa_ok = row[3] is not None and abs(float(row[3]) - exp_oa) < tol
                if ca_ok and fa_ok and oa_ok:
                    found_years.add(yr)
                else:
                    all_values_ok = False
                    break
            except (TypeError, ValueError):
                all_values_ok = False
                break
        if all_values_ok and len(found_years) == 5:
            score += 0.5
    if len(data_rows) >= 5:
        try:
            actual_years = [int(row[0]) for row in data_rows[:5] if row[0] is not None]
            if actual_years == expected_sorted_years:
                score += 0.5
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_xlsx_padded_count__50235d7deccc7f57d815dc66d0e983bc_qw35sft2_3b6a4b8c(result, expected, **options):
    """Check D column zero-padding (0.5) and E1 customer count (0.5)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    rows = result.get('rows', [])
    if rows:
        correct = sum((1 for r in rows if r.get('d') == str(r['c']).zfill(7)))
        if correct == len(rows):
            score += 0.5
    count_expected = expected.get('count_expected', 29)
    e1 = result.get('e1')
    try:
        if e1 is not None and int(e1) == count_expected:
            score += 0.5
    except (ValueError, TypeError):
        pass
    return score

def xlsx_spent_format_total_row__40a8f46c1e43b9d32b31b2230ffda8db_qw35sft2_4e899856(result, expected, **options):
    """Check C2:C8 have 2dp format, A9 = 'Total', C9 = sum of spent values."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    c_formats = result.get('c_formats', [])
    has_two_dp = bool(c_formats) and all((fmt is not None and ('0.00' in str(fmt) or '##0.00' in str(fmt)) for fmt in c_formats))
    if has_two_dp:
        score += 0.34
    a9_value = result.get('a9_value')
    expected_label = expected.get('expected_label', 'Total')
    if a9_value is not None and str(a9_value).strip().lower() == expected_label.lower():
        score += 0.33
    c9_val = result.get('c9_value')
    expected_sum = expected.get('expected_sum', 288.09)
    if c9_val is not None:
        try:
            if abs(float(c9_val) - expected_sum) < 0.01:
                score += 0.33
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_sheet_rename_and_pdf__57d2609933a3290edf4943a154f0921f_qw35sft2_c493ce14(result, expected, **options):
    """Partial credit: 0.33 for correct sheet rename, 0.34 for fit-to-one-page scaling, 0.33 for PDF exported."""
    score = 0.0
    expected_name = expected.get('sheet_name', 'Attendance')
    if result.get('sheet_name') == expected_name:
        score += 0.33
    if result.get('page_fit', False):
        score += 0.34
    if result.get('pdf_exists', False):
        score += 0.33
    return score

def check_calc_pivot_and_total__2c779926d5355737435372c9e97f0d18_qw35sft2_cdf14c34(result, expected, **options):
    """
    Partial credit:
      0.5 - Sheet2 exists with correct pivot table (invoice 10505 count == 5)
      0.5 - Sheet1.G20 == 7071 (SUM of Sales column G2:G19)
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    sheet2_exists = result.get('sheet2_exists', False)
    count_10505 = result.get('invoice_count_10505')
    if sheet2_exists:
        try:
            if int(count_10505) == 5:
                score += 0.5
            else:
                score += 0.25
        except (TypeError, ValueError):
            score += 0.25
    expected_total = expected.get('total_sales', 7071)
    sheet1_g20 = result.get('sheet1_g20')
    try:
        if abs(float(sheet1_g20) - float(expected_total)) < 0.01:
            score += 0.5
    except (TypeError, ValueError):
        pass
    return min(score, 1.0)

def check_pivot_and_sorted__3da7df8645c1252be37e4e927311c6ef_qw35sft2_6755b7e0(result, expected, **options):
    """Partial credit: 0.5 for Sheet2 having all promotion revenues, 0.5 for Sheet1 sorted by Revenue descending."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    tolerance = expected.get('tolerance', 1.0)
    score = 0.0
    expected_revenues = [float(r) for r in expected.get('expected_revenues', [])]
    sheet2_numeric = result.get('sheet2_numeric', [])
    if result.get('sheet2_exists') and expected_revenues:
        found = sum((1 for exp_rev in expected_revenues if any((abs(v - exp_rev) <= tolerance for v in sheet2_numeric))))
        score += 0.5 * (found / len(expected_revenues))
    top_revenues = result.get('sheet1_top_revenues', [])
    max_expected = float(expected.get('max_revenue', 1629.58))
    if top_revenues and abs(top_revenues[0] - max_expected) <= tolerance:
        score += 0.5
    elif len(top_revenues) >= 2 and top_revenues[0] >= top_revenues[1]:
        score += 0.25
    return min(score, 1.0)

def check_cell_approx__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_4b449af3(result, expected, **options):
    """Check that Sheet1!I1 holds the expected numeric value within tolerance."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    if actual is None:
        return 0.0
    expected_val = expected.get('value')
    tolerance = expected.get('tolerance', 0.1)
    try:
        return 1.0 if abs(float(actual) - float(expected_val)) <= tolerance else 0.0
    except (TypeError, ValueError):
        return 0.0

def check_pivot_and_sum__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_700872ff(result, expected, **options):
    """Partial-scoring metric for three sub-goals:
      0.4 - Sheet2 contains a product pivot table (detected via 'product' text)
      0.4 - Sheet2 contains a sales-channel pivot table (detected via 'channel' text)
      0.2 - Sheet1!I1 SUM value matches expected within tolerance
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    if result.get('sheet2_exists') and result.get('sheet2_has_product'):
        score += 0.4
    if result.get('sheet2_exists') and result.get('sheet2_has_channel'):
        score += 0.4
    actual_i1 = result.get('sheet1_i1')
    if actual_i1 is not None:
        expected_val = expected.get('value', 2368.15)
        tolerance = expected.get('tolerance', 0.1)
        try:
            if abs(float(actual_i1) - float(expected_val)) <= tolerance:
                score += 0.2
        except (TypeError, ValueError):
            pass
    return min(score, 1.0)

def check_sheet_exists__0ed76c34ce42b49c8d0f69a316afbc19_qw35sft2_e4d74bcc(result, expected, **options):
    """Check if a sheet with the expected name exists in the workbook."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    sheet_name = expected.get('sheet_name', 'Sheet2')
    return 1.0 if sheet_name in result.get('sheet_names', []) else 0.0

def check_cell_text__fd1c3b4eab4669bd9848fe12e4f0d1d4_qw35sft2_9cce5d4b(result, expected, **options):
    """Check if a cell value matches expected text (case-insensitive strip)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value', '')
    if actual is None:
        return 0.0
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_cells_ordered_values__1d24482d98b418ee12818e4d70230d5c_qw35sft2_5772c294(result, expected, **options):
    """Check that all 6 unique names are entered in the correct order, with partial credit per entry."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    actual_values = result.get('values', [])
    expected_values = expected.get('expected_values', [])
    if not expected_values:
        return 0.0
    n = len(expected_values)
    correct = 0
    for i, exp_val in enumerate(expected_values):
        if i < len(actual_values):
            actual = actual_values[i]
            if actual is not None and str(actual).strip() == str(exp_val).strip():
                correct += 1
    return correct / n

def check_sheet_exists__08c531e403541f6bab4f59a478d0e6c2_qw35sft2_69a71ad6(result, expected, **options):
    """Check if a specific sheet name exists in the workbook."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    sheet_names = result.get('sheet_names', [])
    expected_sheet = expected.get('expected_sheet', 'Sheet2')
    return 1.0 if expected_sheet in sheet_names else 0.0

def check_sheet_and_cell__ebd5692a1e5abf6c2058f3bfb502b3de_qw35sft2_24eddc65(result, expected, **options):
    """Partial-credit check for sheet renames plus cell A1 update in backup sheet.

    Scoring:
    - 0.50: sheet names are correct and in order
    - 0.50: cell A1 of 'LARS Resources (Backup)' matches expected value
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    actual_sheets = result.get('sheet_names', [])
    expected_sheets = expected.get('sheet_names', [])
    if actual_sheets == expected_sheets:
        score += 0.5
    actual_a1 = result.get('backup_a1', '')
    expected_a1 = expected.get('backup_a1', '')
    if actual_a1 and expected_a1 and (str(actual_a1).strip() == str(expected_a1).strip()):
        score += 0.5
    return min(score, 1.0)

def check_calc_sheet2_3col__581360a1329c03f025e39497ca2b0766_qw35sft2_1a12ad34(result, expected, **options):
    """Check Sheet2 has 3-column headers and correct sum values. Partial credit per cell."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    score = 0.0
    weight = 1.0 / 6
    checks = [('A1', expected.get('header_revenue', 'Total Revenue')), ('B1', expected.get('header_expenses', 'Total Expenses')), ('C1', expected.get('header_net_income', 'Net Income')), ('A2', expected.get('sum_revenue', 867786)), ('B2', expected.get('sum_expenses', 411686)), ('C2', expected.get('sum_net_income', 456100))]
    for key, exp_val in checks:
        raw = result.get(key)
        actual = int(raw) if isinstance(raw, float) and raw == int(raw) else raw
        if isinstance(exp_val, str):
            if isinstance(actual, str) and actual.strip() == exp_val.strip():
                score += weight
        elif actual == exp_val:
            score += weight
    return round(min(score, 1.0), 4)

def check_calc_earned_and_hours__ddba4585b2adadf3e6cdf7efce56b822_qw35sft2_e04c2448(result, expected, **options):
    """Partial credit: 0.5 for correct E3 (total earned ~191.67), 0.5 for correct G3 (decimal hours ~7.667)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    tolerance = expected.get('tolerance', 0.01)
    e3 = result.get('e3_value')
    expected_e3 = expected.get('expected_e3', 191.6667)
    if e3 is not None and abs(float(e3) - float(expected_e3)) <= tolerance:
        score += 0.5
    g3 = result.get('g3_value')
    expected_g3 = expected.get('expected_g3', 7.6667)
    if g3 is not None and abs(float(g3) - float(expected_g3)) <= tolerance:
        score += 0.5
    return score

def check_cells_red__cc3c6531bd3eb1e9b2189f82dea73399_qw35sft2_e0fd9222(result, expected, **options):
    """Check that all target cells have the expected red background color.
    Returns partial credit based on fraction of target cells correctly colored."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    target_cells = expected.get('target_cells', [])
    expected_color = expected.get('expected_color', 'FFFF0000').upper()
    if not target_cells:
        return 0.0
    correct = sum((1 for cell in target_cells if result.get(cell, '').upper() == expected_color))
    return correct / len(target_cells)

def check_cell_value_equals__9d8210ef208924a20f5946d708d3b9b8_qw35sft2_a6259fb1(result, expected, **options):
    """Return 1.0 if the cell value matches expected_value (numeric or string)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else result
    expected_val = expected.get('expected_value')
    if actual is None:
        return 0.0
    if isinstance(expected_val, (int, float)) and isinstance(actual, (int, float)):
        return 1.0 if abs(float(actual) - float(expected_val)) < 0.01 else 0.0
    if isinstance(expected_val, str) and isinstance(actual, str):
        return 1.0 if actual.strip().lower() == expected_val.strip().lower() else 0.0
    return 1.0 if actual == expected_val else 0.0

def check_pptx_table_cells__d704f09ae59c96034fafb3c1ef256369_qw35sft2_e6a9ffc5(result, expected, **options):
    """Check first row headers (0.7) and a specific cell value (0.3). Partial credit."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    expected_row = expected.get('expected_table_row0', [])
    actual_row = result.get('table_row0', [])
    if expected_row and actual_row and (len(actual_row) == len(expected_row)):
        if all((a == e for a, e in zip(actual_row, expected_row))):
            score += 0.7
    expected_cell = expected.get('expected_cell', '')
    actual_cell = result.get('specific_cell', '')
    if expected_cell and actual_cell == expected_cell:
        score += 0.3
    return min(score, 1.0)

def check_table_7x5_last_cell__818c616cbd97d0d5bf2ff81918416501_qw35sft2_1a57fa14(result, expected, **options):
    """Check 7x5 table inserted and last cell of the table has expected text.
    Partial credit: 0.5 for correct table/dimensions, 0.5 for last cell content.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count') == expected.get('table_count', 2) and result.get('last_rows') == expected.get('last_rows', 5) and (result.get('last_cols') == expected.get('last_cols', 7)):
        score += 0.5
    expected_text = expected.get('last_cell', '')
    actual_text = result.get('last_cell', '')
    if expected_text and expected_text.lower() in actual_text.lower():
        score += 0.5
    return score

def check_table_7x5_with_cell__2aa481b876a15405aedb5c91a8fc78e9_qw35sft2_8d4c3375(result, expected, **options):
    """Check 7x5 table was inserted and cell A1 of the new table has the expected text.
    Partial credit: 0.5 for correct table/dimensions, 0.5 for cell content.
    """
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    score = 0.0
    if result.get('table_count') == expected.get('table_count', 2) and result.get('last_rows') == expected.get('last_rows', 5) and (result.get('last_cols') == expected.get('last_cols', 7)):
        score += 0.5
    expected_text = expected.get('cell_a1', '')
    actual_text = result.get('cell_a1', '')
    if expected_text and expected_text.lower() in actual_text.lower():
        score += 0.5
    return score

def check_xlsx_cell_val__2feaa23241a59b8898b31056b9df1f85_qw35sft2_706eb125(result, expected, **options):
    """Check that a specific cell contains the expected string value (case-insensitive)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    if str(actual).strip().lower() == str(expected_value).strip().lower():
        return 1.0
    return 0.0

def check_file_numeric_content__f36e167c6b956664e01e5e47712b671b_qw35sft2_4e4217aa(result, expected, **options):
    """Check that a text file contains the expected numeric count."""
    if result.get('error') and (not result.get('content')):
        return 0.0
    content = result.get('content', '').strip()
    expected_count = str(expected.get('count', ''))
    return 1.0 if content == expected_count else 0.0

def check_xlsx_freeze_panes__b12f1d02d5521f031b94e1747a1c556a_qw35sft2_3f46aafb(result, expected, **options):
    """Check that freeze_panes is set to freeze the top row (A2)."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_freeze = expected.get('expected_freeze', 'A2')
    actual_freeze = result.get('freeze_panes')
    return 1.0 if actual_freeze == expected_freeze else 0.0

def check_xlsx_year_column__c99ee3ea1307143970f7d9489f51bcb4_qw35sft2_457e30c1(result, expected, **options):
    """Partial credit for Year column: 0.2 for correct header, 0.16 each for 5 correct year values."""
    if result is None or result.get('error'):
        return 0.0
    score = 0.0
    expected_header = expected.get('expected_header', 'Year')
    actual_header = result.get('header')
    if actual_header is not None and str(actual_header).strip().lower() == expected_header.lower():
        score += 0.2
    expected_years = expected.get('expected_years', [2018, 2019, 2017, 2018, 2018])
    actual_years = result.get('years', [])
    for i, exp_yr in enumerate(expected_years):
        if i < len(actual_years) and actual_years[i] is not None:
            try:
                if int(actual_years[i]) == int(exp_yr):
                    score += 0.16
            except (ValueError, TypeError):
                pass
    return min(round(score, 4), 1.0)

def check_xlsx_cell_value__5777a0629a4e670ae915b1fe64e58378_qw35sft2_e7be0bcf(result, expected, **options):
    if not result or result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None or expected_value is None:
        return 0.0
    return 1.0 if str(actual).strip() == str(expected_value).strip() else 0.0

def check_xlsx_cell_val__61d78e261ca58a364d140d5219ceb5c7_qw35sft2_922f8d42(result, expected, **options):
    """Check that a specific cell contains the expected value."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    try:
        if float(actual) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_xlsx_header_align__7d3ab32fb77b1ea64cfabfdcd56b69aa_qw35sft2_c23af72a(result, expected, **options):
    """Check that all header row cells A1:D1 are center-aligned."""
    if not isinstance(result, dict) or result.get('error'):
        return 0.0
    expected_align = expected.get('expected_alignment', 'center')
    alignments = result.get('header_alignments', [])
    if not alignments:
        return 0.0
    return 1.0 if all((a == expected_align for a in alignments)) else 0.0

def check_xlsx_row_data__5abc6535b144407839fbbe3d8a497678_qw35sft2_41d3f63e(result, expected, **options):
    if not result or (isinstance(result, dict) and result.get('error')):
        return 0.0
    score = 0.0
    checks = [k for k in expected if not k.startswith('_')]
    if not checks:
        return 0.0
    per_check = 1.0 / len(checks)
    for key in checks:
        expected_val = expected[key]
        actual_val = result.get(key)
        if actual_val is None:
            continue
        if isinstance(expected_val, (int, float)) and isinstance(actual_val, (int, float)):
            if abs(float(actual_val) - float(expected_val)) < 1:
                score += per_check
        elif str(actual_val).strip().lower() == str(expected_val).strip().lower():
            score += per_check
    return min(score, 1.0)

def check_xlsx_cell_val__9d1e0f5a79d29ecfac0c206e85fa98e9_qw35sft2_73735a2c(result, expected, **options):
    """Check that a specific cell contains the expected value (numeric or string)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value')
    expected_value = expected.get('expected_value')
    if actual is None:
        return 0.0
    try:
        if float(actual) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    if str(actual).strip() == str(expected_value).strip():
        return 1.0
    return 0.0

def check_unseen_sheet_headers__ac461f3f5169eddd647d2f7d2aad85a2_qw35sft2_0281101b(result, expected, **options):
    """Check that the 'unseen_movies' sheet exists and has the correct headers."""
    if not result or result.get('error'):
        return 0.0
    expected_headers = expected.get('expected_headers', ['title', 'release year', 'ratings', 'description'])
    actual_headers = result.get('headers', [])
    if not actual_headers:
        return 0.0
    actual_headers = [h for h in actual_headers if h is not None]
    if actual_headers == expected_headers:
        return 1.0
    match_count = sum((1 for a, e in zip(actual_headers, expected_headers) if a == e))
    return round(match_count / len(expected_headers), 2)

def check_xlsx_cell_value__9d415d4db27ab3e1c21d0be924167f0e_qw35sft2_fc10d580(result, expected, **options):
    """Check if a cell value matches the expected string (case-insensitive)."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    actual = result.get('value') if isinstance(result, dict) else None
    if actual is None:
        return 0.0
    expected_value = expected.get('expected_value', '')
    if actual.strip().lower() == expected_value.strip().lower():
        return 1.0
    return 0.0
