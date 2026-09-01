"""VeriGen RL judge functions.

Source: metrics.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
import base64

logger = logging.getLogger(__name__)
logger_qw35sft2_60bfba = logging.getLogger(__name__)
logger_qw35sft2_82cd1a = logging.getLogger(__name__)
logger_qw35sft2_3433d9 = logging.getLogger(__name__)
logger_qw35sft2_c98520 = logging.getLogger('desktopenv.metrics.gimp_custom')
_ICML_CITY_ALIASES_qw35sft2_45c2e8 = {'new york': ['new york', 'new york city', 'nyc'], 'long beach': ['long beach', 'los angeles', 'la']}
logger_qw35sft2_d0992a = logging.getLogger('desktopenv.metrics.eml_backup')
logger_qw35sft2_fca153 = logging.getLogger(__name__)
logger_qw35sft2_2fd121 = logging.getLogger(__name__)
logger_qw35sft2_dce5f0 = logging.getLogger('desktopenv.metrics.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_34eb84 = logging.getLogger(__name__)
logger_qw35sft2_103ddb = logging.getLogger('desktopenv.metrics.eml_subject')
logger_qw35sft2_1d640f = logging.getLogger(__name__)
logger_qw35sft2_f5fbc6 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_a1dd18 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_9a65d1 = logging.getLogger(__name__)
logger_qw35sft2_2ef5dd = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_592596 = logging.getLogger(__name__)
_DEFAULT_BRIGHTNESS_qw35sft2_fa3a1f = 1.0
_BRIGHTNESS_EPSILON_qw35sft2_fa3a1f = 0.01
logger_qw35sft2_061fea = logging.getLogger(__name__)
logger_qw35sft2_e8a2da = logging.getLogger('desktopenv.metrics.vlc_play_stop')
logger_qw35sft2_462da1 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_a878d7 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_d55a5c = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_8ff4c2 = logging.getLogger(__name__)
logger_qw35sft2_0c2f54 = logging.getLogger(__name__)
logger_qw35sft2_211cdf = logging.getLogger('desktopenv.metrics.vlc_traj_verify_1')
logger_qw35sft2_868f56 = logging.getLogger(__name__)
logger_qw35sft2_2d5b02 = logging.getLogger(__name__)
logger_qw35sft2_1e51cc = logging.getLogger('desktopenv.metrics.vlc_next')
logger_qw35sft2_ccdace = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_544f1c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_1be130 = logging.getLogger('desktopenv.metrics.vlc_traj_verify_4')
logger_qw35sft2_391de6 = logging.getLogger(__name__)
logger_qw35sft2_b4d0e7 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_2c2c74 = logging.getLogger(__name__)
_SATURATION_GRAYSCALE_MAX_qw35sft2_7934da = 0.1
logger_qw35sft2_c58a25 = logging.getLogger(__name__)
logger_qw35sft2_04b5ee = logging.getLogger(__name__)
logger_qw35sft2_9d1c6a = logging.getLogger('desktopenv.metrics.vlc_vol_up')
logger_qw35sft2_a640c9 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_b6f986 = logging.getLogger(__name__)
logger_qw35sft2_b9d146 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_256b9d = logging.getLogger(__name__)
logger_qw35sft2_208a9d = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_99c515 = logging.getLogger(__name__)
logger_qw35sft2_d8f706 = logging.getLogger(__name__)
logger_qw35sft2_0717dc = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_849761 = logging.getLogger('desktopenv.metrics.vlc_play_recording')
logger_qw35sft2_43c1c4 = logging.getLogger(__name__)
_DEFAULT_CONTRAST_qw35sft2_bffd0c = 1.0
_CONTRAST_EPSILON_qw35sft2_bffd0c = 0.01
logger_qw35sft2_813e96 = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_589392 = logging.getLogger(__name__)
logger_qw35sft2_14d58e = logging.getLogger(__name__)
logger_qw35sft2_109219 = logging.getLogger('desktopenv.metrics.vlc_custom')
logger_qw35sft2_55647c = logging.getLogger('desktopenv.metrics.vlc_traj')
logger_qw35sft2_421d15 = logging.getLogger(__name__)
logger_qw35sft2_08d9ad = logging.getLogger('desktopenv.metrics.vlc_play_next')

__all__ = ['check_folder_created__1d547c7422f60e4be92e89ee2379630a', 'check_local_folder__2731b9abd5cfbad9ed4df8aae737addc', 'check_folder_contents__448eb54c6f660c787448d046dde3cea3', 'check_files_in_folder__1a224fb0890014daa398af667e7e35c0', 'check_files_moved_to_folder__2cb8499f32af5ada0bb0e955434e59dd', 'check_file_in_folder__69659eabd73459db19e2b0d1637cc759', 'check_music_folder_move__eed664169e1770df941a91e7673f0e98', 'check_files_in_folder__9aea5fe5a12921799c11d30552224cd9', 'check_desktop_file__bc88c1a76af0c575a8c22634fe96ef8b_qw35sft2_a7d2f710', 'check_file_on_desktop__00bb03d81e57baa40bd1c86ce0b1574d_qw35sft2_df86431d', 'check_file_on_desktop__b0e8ccdfade877a9b91932809683825c_qw35sft2_e97098dd', 'check_vim_removed_terminal_added__f6775b09b9b85d12e5fafc15ee8143de_qw35sft2_0232da1f', 'check_bashrc_python_settings__dd32b078ecc7acd064c9c431fc6fece6_qw35sft2_dd1073ae', 'check_rename_and_subfolder__cc8a90beff8b6d7aa9d9f79938c158e4_qw35sft2_4505b83c', 'check_volume_terminal__bafe7fff6ff48b4d9b5b1fd5f2d7f539_qw35sft2_4155aa0b', 'check_os_night_dark__0de6f297debf99316b6fa7b7c038fa5d_qw35sft2_c0558223', 'check_folder_view_active__d9dcbaa59677e530494c93dfe8659800_qw35sft2_31567fbd', 'check_folder_view_active__0dc2b9b9129cecee63a42fd11a56a3f7_qw35sft2_58e4e9c2', 'check_folder_view_active__3098a0959b134380c33dcaefd08a2ca7_qw35sft2_612b9de0', 'check_multi_folder_views__582d00dc650f37f3e769d70699c5f1b6_qw35sft2_30cd7af6', 'check_multi_folder_views__9ec4687f5770fc31e26559c7787fe14b_qw35sft2_eb9975ee']

def check_folder_created__1d547c7422f60e4be92e89ee2379630a(result, expected, **options):
    """Check if the expected folder was created in Local Folders."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if result.get('folder_exists', False):
        return 1.0
    return 0.0

def check_local_folder__2731b9abd5cfbad9ed4df8aae737addc(result, expected, **options):
    """Check if expected local folder exists in Thunderbird."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    folders = result.get('folders', [])
    expected_folder = expected.get('expected_folder', '')
    if not expected_folder:
        return 0.0
    for folder in folders:
        if folder.lower() == expected_folder.lower():
            return 1.0
    return 0.0

def check_folder_contents__448eb54c6f660c787448d046dde3cea3(result, expected, **options):
    """Check if folder contains expected files. Partial credit per file."""
    if not isinstance(result, dict) or not result.get('exists', False):
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    actual_files = set(result.get('files', []))
    matches = sum((1 for f in expected_files if f in actual_files))
    return matches / len(expected_files)

def check_files_in_folder__1a224fb0890014daa398af667e7e35c0(result, expected, **options):
    """Check if expected files are present in the folder listing with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_files = result.get('files', [])
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    found = sum((1 for f in expected_files if f in actual_files))
    return found / len(expected_files)

def check_files_moved_to_folder__2cb8499f32af5ada0bb0e955434e59dd(result, expected, **options):
    """Check if folder was created and files were moved into it. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    expected_files = expected.get('expected_files', [])
    score = 0.0
    if result.get('folder_exists'):
        score += 0.25
    if expected_files:
        per_file = 0.75 / len(expected_files)
        subfolder_files = result.get('subfolder_files', [])
        for f in expected_files:
            if f in subfolder_files:
                score += per_file
    return min(score, 1.0)

def check_file_in_folder__69659eabd73459db19e2b0d1637cc759(result, expected, **options):
    """Check if expected file is present in the folder listing."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_files = result.get('files', [])
    expected_file = expected.get('expected_file', '')
    if expected_file in actual_files:
        return 1.0
    return 0.0

def check_music_folder_move__eed664169e1770df941a91e7673f0e98(result, expected, **options):
    """Check if Classics dir was created and MP3 files moved. Partial credit."""
    if isinstance(result, str) or result.get('error'):
        return 0.0
    score = 0.0
    expected_files = expected.get('expected_files', [])
    if result.get('dir_exists'):
        score += 0.2
    if expected_files and result.get('files_in_classics'):
        per_file = 0.8 / len(expected_files)
        for f in expected_files:
            if f in result['files_in_classics']:
                score += per_file
    return min(round(score, 2), 1.0)

def check_files_in_folder__9aea5fe5a12921799c11d30552224cd9(result, expected, **options):
    """Check if expected files are present in the folder listing with partial credit."""
    if isinstance(result, str) and 'error' in result.lower():
        return 0.0
    actual_files = result.get('files', [])
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    found = sum((1 for f in expected_files if f in actual_files))
    return found / len(expected_files)

def check_desktop_file__bc88c1a76af0c575a8c22634fe96ef8b_qw35sft2_a7d2f710(result, expected, **options):
    """Check if the specific Desktop file exists (binary 0.0 or 1.0)."""
    if not result or not isinstance(result, dict):
        return 0.0
    if result.get('error'):
        return 0.0
    return 1.0 if result.get('exists', False) else 0.0

def check_file_on_desktop__00bb03d81e57baa40bd1c86ce0b1574d_qw35sft2_df86431d(result, expected, **options):
    """Check if a specific filename appears in the Desktop listing."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    filename = expected.get('filename', '')
    if filename and filename in output:
        return 1.0
    return 0.0

def check_file_on_desktop__b0e8ccdfade877a9b91932809683825c_qw35sft2_e97098dd(result, expected, **options):
    """Check if a specific filename appears in the Desktop listing."""
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    filename = expected.get('filename', '')
    if filename and filename in output:
        return 1.0
    return 0.0

def check_vim_removed_terminal_added__f6775b09b9b85d12e5fafc15ee8143de_qw35sft2_0232da1f(result, expected, **options):
    """
    Partial-credit metric for: remove vim.desktop AND add org.gnome.Terminal.desktop.
    - 0.5 for vim.desktop absent from favorites
    - 0.5 for add_app (from expected rules) present in favorites
    Returns float in [0.0, 1.0].
    """
    if isinstance(result, dict) and result.get('error'):
        return 0.0
    apps = result.get('apps', [])
    if not isinstance(apps, list):
        return 0.0
    score = 0.0
    if 'vim.desktop' not in apps:
        score += 0.5
    add_app = expected.get('add_app', '')
    if add_app and add_app in apps:
        score += 0.5
    return score

def check_bashrc_python_settings__dd32b078ecc7acd064c9c431fc6fece6_qw35sft2_dd1073ae(result, expected, **options):
    """
    Partial-credit check for two ~/.bashrc Python settings:
      - 0.5 if 'alias py' is present (py alias for python3)
      - 0.5 if 'PYTHONDONTWRITEBYTECODE' is present
    """
    if not result or not isinstance(result, dict):
        return 0.0
    content = result.get('bashrc_content', '')
    score = 0.0
    if 'alias py' in content:
        score += 0.5
    if 'PYTHONDONTWRITEBYTECODE' in content:
        score += 0.5
    return score

def check_rename_and_subfolder__cc8a90beff8b6d7aa9d9f79938c158e4_qw35sft2_4505b83c(result, expected, **options):
    """
    Partial-credit check:
      0.5 - Desktop folder renamed to todo_list_Jan_2
      0.5 - 'tasks' subfolder created inside todo_list_Jan_2
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('renamed'):
        score += 0.5
    if result.get('subfolder'):
        score += 0.5
    return score

def check_volume_terminal__bafe7fff6ff48b4d9b5b1fd5f2d7f539_qw35sft2_4155aa0b(result, expected, **options):
    """Partial credit: 0.5 for max volume (100%), 0.5 for terminal window open."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if '100%' in result.get('volume_output', ''):
        score += 0.5
    if result.get('terminal_open', False):
        score += 0.5
    return score

def check_os_night_dark__0de6f297debf99316b6fa7b7c038fa5d_qw35sft2_c0558223(result, expected, **options):
    """Check that Night Light is enabled (0.5) and dark mode is active (0.5).

    expected (dict, already unwrapped from rules by get_rule):
      {
        "night_light_enabled": "true",
        "color_scheme": "prefer-dark"
      }
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    night_actual = result.get('night_light_enabled', '')
    night_expected = expected.get('night_light_enabled', 'true')
    if night_expected in night_actual:
        score += 0.5
    dark_actual = result.get('color_scheme', '')
    dark_expected = expected.get('color_scheme', 'prefer-dark')
    if dark_expected in dark_actual:
        score += 0.5
    return score

def check_folder_view_active__d9dcbaa59677e530494c93dfe8659800_qw35sft2_31567fbd(result, expected, **options):
    """Check that a specific folder view mode is active in Thunderbird."""
    if result.get('error') or result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    required_view = expected.get('required_view', '')
    if required_view and required_view in mode:
        return 1.0
    return 0.0

def check_folder_view_active__0dc2b9b9129cecee63a42fd11a56a3f7_qw35sft2_58e4e9c2(result, expected, **options):
    """Check that a specific folder view mode is active in Thunderbird."""
    if result.get('error') or result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    required_view = expected.get('required_view', '')
    if required_view and required_view in mode:
        return 1.0
    return 0.0

def check_folder_view_active__3098a0959b134380c33dcaefd08a2ca7_qw35sft2_612b9de0(result, expected, **options):
    """Check that a specific folder view mode is active in Thunderbird."""
    if result.get('error') or result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    required_view = expected.get('required_view', '')
    if required_view and required_view in mode:
        return 1.0
    return 0.0

def check_multi_folder_views__582d00dc650f37f3e769d70699c5f1b6_qw35sft2_30cd7af6(result, expected, **options):
    """Check that multiple folder view modes are active in Thunderbird (partial credit per view)."""
    if result.get('error') or result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    required_views = expected.get('required_views', [])
    if not required_views:
        return 0.0
    score = 0.0
    per_view = 1.0 / len(required_views)
    for view in required_views:
        if view in mode:
            score += per_view
    return min(score, 1.0)

def check_multi_folder_views__9ec4687f5770fc31e26559c7787fe14b_qw35sft2_eb9975ee(result, expected, **options):
    """Check that multiple folder view modes are active in Thunderbird (partial credit per view)."""
    if result.get('error') or result.get('mode') is None:
        return 0.0
    mode = result.get('mode', '')
    required_views = expected.get('required_views', [])
    if not required_views:
        return 0.0
    score = 0.0
    per_view = 1.0 / len(required_views)
    for view in required_views:
        if view in mode:
            score += per_view
    return min(score, 1.0)
