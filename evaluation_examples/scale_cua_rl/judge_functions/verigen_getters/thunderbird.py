"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_thunderbird_smtp_settings__22e52b73b299835b6ec6765d839ed989', 'get_thunderbird_eml_backup__b0f7174902a45f1ac3121d8b78a62ff1', 'get_thunderbird_inbox_state__8e4a0a40e4d045486197a43a70e9b321', 'get_tb_local_folder_list__4a5755b15fd8848dbc70d01e3dbb4e89', 'get_thunderbird_local_folder__4de43750fd473996ecf007ebf44296c8', 'get_thunderbird_filter_file__9b422e9b7b3e9e70da3aceccdee5065b', 'get_thunderbird_folder_state__168d8b8e9196711d91fc3322429f68b5', 'get_thunderbird_smtp_settings__2c49114272e52893d5a44b081f8b14e0', 'get_thunderbird_account_settings__20583839ff43876dc06848262d06ce9a', 'get_thunderbird_local_folder__2731b9abd5cfbad9ed4df8aae737addc', 'get_tb_bills_folder_content__d9f2b609b5f34a856608f932bab06223', 'get_tb_daily_msg_count__c7401fcb5b9f30e5bd0dd82a92e3f2df', 'get_thunderbird_identity__4a5ace1e770bd482f5a0433cd05f4404', 'get_git_email_and_push__0c30b3c3ecef4bb03c9efc8f2c5c6e42_qw35sft2_cb5cf810', 'get_thunderbird_folder_view__d9dcbaa59677e530494c93dfe8659800_qw35sft2_d8ca49d4', 'get_thunderbird_multi_prefs__840f3ea7ded7f1b6695bc630a0a035b7_qw35sft2_ee0362cb', 'get_thunderbird_local_folders__1ea15ae7572bc843dca8cde6cbcd3a7c_qw35sft2_8c8fdc3d', 'get_tb_forward_filter_state__f7fbe032caaa49b6e5ebdeafec739cf8_qw35sft2_916711f4', 'get_thunderbird_smtp_no_incoming__f562ce0912b384f5748be41ba253880a_qw35sft2_0a0eb666', 'get_tb_acct_folder__2e9b3985bbecaa462dbe43eedf3d22fe_qw35sft2_479e0c2d', 'get_thunderbird_draft_cc_attachment__ae64d5efe053bb1ed61ef261c2f7eddd_qw35sft2_b04c87df', 'get_thunderbird_imap_setup_state__7fc0ec566fbfa7a7ef758884ad90b95b_qw35sft2_02f51ec2', 'get_tb_folder_two_filters__bc08452635875f98955858c473a86873_qw35sft2_574f5cb0', 'get_thunderbird_folder_view__0dc2b9b9129cecee63a42fd11a56a3f7_qw35sft2_342c6d1b', 'get_thunderbird_nested_folders__1d59e1d1688abeb99fc9fce073e49b55_qw35sft2_f9dee36b', 'get_tb_acct_smtp__93d97cb63d509700a5843e46760d10d0_qw35sft2_684ba59d', 'get_thunderbird_smtp_full__4c184097a04ef7d98ec36b30ff774f5c_qw35sft2_2427a295', 'get_thunderbird_filter_names__e046edd8a7ea67608e0814ad68f68ef0_qw35sft2_adcbd111', 'get_thunderbird_compose_subject_attachment__59182577b48daf97ea5874cad0371f93_qw35sft2_1360ba35', 'get_tb_dual_filters__f4464353c01b135bac464c93385b0943_qw35sft2_42628857', 'get_thunderbird_manual_cfg__598710d21fbdbbb5b9d697d62aaf5b30_qw35sft2_ae05246b', 'get_tb_three_folders_one_filter__9865f0d026c5662c0f08b728c450437f_qw35sft2_3d152545', 'get_thunderbird_multi_prefs__f88a527105940879c254b88e4c3fea68_qw35sft2_8c851feb', 'get_thunderbird_folder_view__3098a0959b134380c33dcaefd08a2ca7_qw35sft2_63010d62', 'get_thunderbird_folder_state__50e2868517db7711c169fe6c468e7766_qw35sft2_9f561ac4', 'get_tb_acct_removal__3f1fce882b4fa7de95c0c939ed3f5b9d_qw35sft2_d7d94b11', 'get_thunderbird_filter_match_all__f5a0e9992343a59dd61c9a37c78c3db0_qw35sft2_4b79d3bf', 'get_thunderbird_draft_attachment__ac90f87152777ec474b581760325d43f_qw35sft2_2768cbc8', 'get_tb_incoming_forward__0ab002743507f6ef7d004b0ee40adbf4_qw35sft2_32a29ccc', 'get_thunderbird_smtp_config__efa5fdf04c5026bfc6f8d6ffd452c5ec_qw35sft2_10c35c16', 'get_thunderbird_imap_port__c8735abb904c2556124a9fb1090745a7_qw35sft2_505ee3a8', 'get_tb_two_folders_two_filters__85999e7d7f538cb3e1b971a9ebe1ed0f_qw35sft2_a5976dbb', 'get_thunderbird_local_folders__e760cb4cf85af456bb04229f8025e52c_qw35sft2_3c45b8b4', 'get_thunderbird_folder_view__582d00dc650f37f3e769d70699c5f1b6_qw35sft2_5b54ad33', 'get_tb_acct_trash__8fc7e10e457380f7cb927cd320b9664e_qw35sft2_66bd48ec', 'get_thunderbird_filter_and_pref__709b71706c15c6652aa94ca6ff4b6ae6_qw35sft2_34b2256c', 'get_tb_matchall_forward__bd4f2f9ff4f1fb7ec8f18ecd69248d51_qw35sft2_8473ec06', 'get_thunderbird_name_email__2d78e59033d5901e822a6624aeea3bc7_qw35sft2_60bf01ac', 'get_thunderbird_compose_cc_attachment__433cf4a73d1a380efac6f447ba251498_qw35sft2_a2b5d268', 'get_thunderbird_smtp_description__8a8a5b00b6516b6e1171569c19648a1b_qw35sft2_cb59c0e2', 'get_tb_two_folders_and_filter__951712fea58d5a6835d2f3a50fd315b8_qw35sft2_cb18f390', 'get_tb_theme_and_folder__bdfb5942d44e1763e7f54cb6f6b53816_qw35sft2_217de1fc', 'get_thunderbird_folder_view__9ec4687f5770fc31e26559c7787fe14b_qw35sft2_8b914604', 'get_thunderbird_filter_incoming__4d2f365259c66d4c891bf0c4c8a5a5f0_qw35sft2_5bcea5d8', 'get_tb_acct_filter__ddc543b6ee27062796019a5514b0e693_qw35sft2_c9574301', 'get_tb_named_forward_filter__e2cb27642e2d718992a0ab3ec3b240f1_qw35sft2_56f302c1', 'get_thunderbird_compose_bcc_attachment__d67afa1e0d57a9bbcae95982d894cc83_qw35sft2_644cc038', 'get_thunderbird_smtp_security__9cc78735939880d6658a6582ff470dcd_qw35sft2_e2f5e05f', 'get_tb_two_folders_two_filters_v2__627b4b586c777792c6f3d931dd087820_qw35sft2_cf0ddccb', 'get_ext_and_minimap__8a991aa9f9913ccca00f4b1c76aac764_qw35sft2_88c50775']

def get_thunderbird_smtp_settings__22e52b73b299835b6ec6765d839ed989(env, config: dict):
    """Get SMTP server settings from Thunderbird prefs.js."""
    prefs_path = config.get('path', '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js')
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='ignore')
    prefs = {}
    for match in re.finditer('user_pref\\("([^"]+)",\\s*(.+?)\\);', content):
        (key, value) = (match.group(1), match.group(2).strip())
        if value.startswith('"') and value.endswith('"'):
            value = value[1:-1]
        elif value == 'true':
            value = True
        elif value == 'false':
            value = False
        else:
            try:
                value = int(value)
            except ValueError:
                try:
                    value = float(value)
                except ValueError:
                    pass
        prefs[key] = value
    smtp_servers = {}
    for (key, value) in prefs.items():
        if key.startswith('mail.smtpserver.') and key != 'mail.smtpserver.default':
            parts = key.split('.')
            if len(parts) >= 4:
                server_id = parts[2]
                setting = '.'.join(parts[3:])
                if server_id not in smtp_servers:
                    smtp_servers[server_id] = {}
                smtp_servers[server_id][setting] = value
    return {'smtp_servers': smtp_servers}

def get_thunderbird_eml_backup__b0f7174902a45f1ac3121d8b78a62ff1(env, config: dict):
    """Check if EML files exist in the specified backup directory."""
    import re
    target_dir = config.get('path', '/home/user/Desktop/inbox_backup')
    result = env.controller.run_bash_script(f'ls -la {target_dir}/*.eml 2>/dev/null', timeout=30)
    stdout = result.get('output', '') if isinstance(result, dict) else str(result)
    eml_files = []
    for line in stdout.strip().split('\n'):
        line = line.strip()
        if line.endswith('.eml'):
            parts = line.split()
            if parts:
                eml_files.append(parts[-1].split('/')[-1])
    return {'eml_files': eml_files, 'count': len(eml_files), 'raw_output': stdout}

def get_thunderbird_inbox_state__8e4a0a40e4d045486197a43a70e9b321(env, config: dict):
    """Check the state of Thunderbird inbox and trash folders."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird')
    result = env.controller.run_bash_script(f'find {profile_path} -maxdepth 1 -type d -name "*.default*" 2>/dev/null | head -1', timeout=30)
    profile_dir = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    inbox_content = ''
    trash_content = ''
    if profile_dir:
        inbox_result = env.controller.run_bash_script(f'find "{profile_dir}/Mail" -name "INBOX" -o -name "Inbox" 2>/dev/null | head -1', timeout=30)
        inbox_path = inbox_result.get('output', '').strip() if isinstance(inbox_result, dict) else str(inbox_result).strip()
        if inbox_path:
            content_result = env.controller.run_bash_script(f'grep -c "^From " "{inbox_path}" 2>/dev/null || echo "0"', timeout=30)
            inbox_content = content_result.get('output', '0').strip() if isinstance(content_result, dict) else str(content_result).strip()
        trash_result = env.controller.run_bash_script(f'find "{profile_dir}/Mail" -name "Trash" 2>/dev/null | head -1', timeout=30)
        trash_path = trash_result.get('output', '').strip() if isinstance(trash_result, dict) else str(trash_result).strip()
        if trash_path:
            content_result = env.controller.run_bash_script(f'grep -c "^From " "{trash_path}" 2>/dev/null || echo "0"', timeout=30)
            trash_content = content_result.get('output', '0').strip() if isinstance(content_result, dict) else str(content_result).strip()
    try:
        inbox_count = int(inbox_content) if inbox_content else 0
    except ValueError:
        inbox_count = 0
    try:
        trash_count = int(trash_content) if trash_content else 0
    except ValueError:
        trash_count = 0
    return {'inbox_message_count': inbox_count, 'trash_message_count': trash_count}

def get_tb_local_folder_list__4a5755b15fd8848dbc70d01e3dbb4e89(env, config: dict):
    """List files in Thunderbird Local Folders directory to check for folder existence."""
    result = env.controller.run_bash_script('PROFILE=$(find /home/user/.thunderbird -maxdepth 1 -name "*.default-release" -type d | head -1) && ls "$PROFILE/Mail/Local Folders/" 2>/dev/null', timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    return {'listing': output}

def get_thunderbird_local_folder__4de43750fd473996ecf007ebf44296c8(env, config: dict):
    """List files in Thunderbird local folders directory."""
    folder_path = config.get('folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/')
    result = env.controller.run_bash_script(f'ls -1 "{folder_path}" 2>/dev/null', timeout=30)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '')
    elif isinstance(result, str):
        output = result
    items = [f.strip() for f in output.split('\n') if f.strip()]
    return {'items': items}

def get_thunderbird_filter_file__9b422e9b7b3e9e70da3aceccdee5065b(env, config: dict):
    """Find and download the Thunderbird msgFilterRules.dat file from the VM."""
    profile_dir = config.get('profile_dir', '/home/user/.thunderbird/t5q2a5hp.default-release/')
    find_cmd = f"find '{profile_dir}' -name 'msgFilterRules.dat' -type f 2>/dev/null | head -1"
    result = env.controller.run_bash_script(find_cmd, timeout=10)
    filter_path = ''
    if isinstance(result, dict):
        filter_path = result.get('output', '').strip()
    elif isinstance(result, str):
        filter_path = result.strip()
    if not filter_path:
        return None
    file_bytes = env.controller.get_file(filter_path)
    if not file_bytes:
        return None
    tmp_path = os.path.join(tempfile.gettempdir(), 'msgFilterRules_v1.dat')
    with open(tmp_path, 'wb') as f:
        f.write(file_bytes)
    return tmp_path

def get_thunderbird_folder_state__168d8b8e9196711d91fc3322429f68b5(env, config: dict):
    """List files in Thunderbird local folders directory for deletion check."""
    folder_path = config.get('folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/')
    result = env.controller.run_bash_script(f'ls -1 "{folder_path}" 2>/dev/null', timeout=30)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '')
    elif isinstance(result, str):
        output = result
    items = [f.strip() for f in output.split('\n') if f.strip()]
    return {'items': items}

def get_thunderbird_smtp_settings__2c49114272e52893d5a44b081f8b14e0(env, config: dict):
    """Get SMTP server settings from Thunderbird prefs.js."""
    prefs_path = config.get('path', '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js')
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='ignore')
    prefs = {}
    for match in re.finditer('user_pref\\("([^"]+)",\\s*(.+?)\\);', content):
        (key, value) = (match.group(1), match.group(2).strip())
        if value.startswith('"') and value.endswith('"'):
            value = value[1:-1]
        elif value == 'true':
            value = True
        elif value == 'false':
            value = False
        else:
            try:
                value = int(value)
            except ValueError:
                try:
                    value = float(value)
                except ValueError:
                    pass
        prefs[key] = value
    smtp_servers = {}
    for (key, value) in prefs.items():
        if key.startswith('mail.smtpserver.') and key != 'mail.smtpserver.default':
            parts = key.split('.')
            if len(parts) >= 4:
                server_id = parts[2]
                setting = '.'.join(parts[3:])
                if server_id not in smtp_servers:
                    smtp_servers[server_id] = {}
                smtp_servers[server_id][setting] = value
    return {'smtp_servers': smtp_servers}

def get_thunderbird_account_settings__20583839ff43876dc06848262d06ce9a(env, config: dict):
    """Get email account and server settings from Thunderbird prefs.js."""
    prefs_path = config.get('path', '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js')
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='ignore')
    prefs = {}
    for match in re.finditer('user_pref\\("([^"]+)",\\s*(.+?)\\);', content):
        (key, value) = (match.group(1), match.group(2).strip())
        if value.startswith('"') and value.endswith('"'):
            value = value[1:-1]
        elif value == 'true':
            value = True
        elif value == 'false':
            value = False
        else:
            try:
                value = int(value)
            except ValueError:
                try:
                    value = float(value)
                except ValueError:
                    pass
        prefs[key] = value
    identities = {}
    for (key, value) in prefs.items():
        if key.startswith('mail.identity.'):
            parts = key.split('.')
            if len(parts) >= 4:
                id_name = parts[2]
                setting = '.'.join(parts[3:])
                if id_name not in identities:
                    identities[id_name] = {}
                identities[id_name][setting] = value
    servers = {}
    for (key, value) in prefs.items():
        if key.startswith('mail.server.') and (not key.startswith('mail.server.default')):
            parts = key.split('.')
            if len(parts) >= 4:
                server_name = parts[2]
                setting = '.'.join(parts[3:])
                if server_name not in servers:
                    servers[server_name] = {}
                servers[server_name][setting] = value
    return {'identities': identities, 'servers': servers}

def get_thunderbird_local_folder__2731b9abd5cfbad9ed4df8aae737addc(env, config: dict):
    """Check if a specific local folder exists in Thunderbird's Local Folders."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird')
    result = env.controller.run_bash_script(f'find {profile_path} -maxdepth 1 -type d -name "*.default*" 2>/dev/null | head -1', timeout=30)
    profile_dir = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    if not profile_dir:
        result = env.controller.run_bash_script(f'find {profile_path} -path "*/Mail/Local Folders" -type d 2>/dev/null | head -1', timeout=30)
        local_folders_path = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    else:
        local_folders_path = f'{profile_dir}/Mail/Local Folders'
    if not local_folders_path:
        return {'error': 'Could not find Local Folders directory', 'folders': []}
    result = env.controller.run_bash_script(f'ls -1 "{local_folders_path}" 2>/dev/null', timeout=30)
    stdout = result.get('output', '') if isinstance(result, dict) else str(result)
    folders = []
    for line in stdout.strip().split('\n'):
        name = line.strip()
        if name and (not name.endswith('.msf')) and (not name.endswith('.sbd')):
            folders.append(name)
    return {'folders': folders, 'local_folders_path': local_folders_path}

def get_tb_bills_folder_content__d9f2b609b5f34a856608f932bab06223(env, config: dict):
    """Check Bills folder mbox for a specific email subject and daily folder message count."""
    script = 'PROFILE=$(find /home/user/.thunderbird -maxdepth 1 -name "*.default-release" -type d | head -1)\nDAILY_COUNT=$(grep -c "^From " "$PROFILE/Mail/Local Folders/daily" 2>/dev/null || echo "0")\nBILLS_HAS_MSG=$(grep -c "HKU Daily News Digest 0124" "$PROFILE/Mail/Local Folders/Bills" 2>/dev/null || echo "0")\necho "daily_count=$DAILY_COUNT"\necho "bills_has_msg=$BILLS_HAS_MSG"'
    result = env.controller.run_bash_script(script, timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    daily_count = -1
    bills_has_msg = 0
    for line in output.strip().split('\n'):
        if line.startswith('daily_count='):
            try:
                daily_count = int(line.split('=')[1])
            except (ValueError, IndexError):
                pass
        if line.startswith('bills_has_msg='):
            try:
                bills_has_msg = int(line.split('=')[1])
            except (ValueError, IndexError):
                pass
    return {'daily_count': daily_count, 'bills_has_msg': bills_has_msg}

def get_tb_daily_msg_count__c7401fcb5b9f30e5bd0dd82a92e3f2df(env, config: dict):
    """Count messages in the daily mbox file by counting 'From ' separator lines."""
    result = env.controller.run_bash_script('PROFILE=$(find /home/user/.thunderbird -maxdepth 1 -name "*.default-release" -type d | head -1) && grep -c "^From " "$PROFILE/Mail/Local Folders/daily" 2>/dev/null || echo "0"', timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    try:
        count = int(output.strip().split('\n')[-1])
    except (ValueError, IndexError):
        count = -1
    return {'count': count}

def get_thunderbird_identity__4a5ace1e770bd482f5a0433cd05f4404(env, config: dict):
    """Get Thunderbird identity setting from prefs.js."""
    prefs_path = config.get('path', '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js')
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    pref_key = config.get('pref_key', 'mail.identity.id1.fullName')
    pattern = 'user_pref\\("' + re.escape(pref_key) + '",\\s*"([^"]*)"\\);'
    match = re.search(pattern, content)
    if match:
        return {'value': match.group(1)}
    else:
        return {'value': None, 'note': f'pref {pref_key} not found in prefs.js'}

def get_git_email_and_push__0c30b3c3ecef4bb03c9efc8f2c5c6e42_qw35sft2_cb5cf810(env, config: dict):
    """
    Get the local git user.email in the binder project and the latest commit message
    in the remote repo.
    Returns dict with keys 'user_email' and 'remote_log'.
    """
    email_result = env.controller.run_bash_script('git -C /home/user/projects/binder config --local user.email 2>&1 || echo ""', timeout=30)
    if isinstance(email_result, dict):
        user_email = email_result.get('stdout', '') or email_result.get('output', '') or str(email_result)
    else:
        user_email = str(email_result) if email_result else ''
    remote_log_result = env.controller.run_bash_script('git -C /home/user/projects/remote_project log --oneline -1 2>&1 || echo "GIT_ERROR"', timeout=30)
    if isinstance(remote_log_result, dict):
        remote_log = remote_log_result.get('stdout', '') or remote_log_result.get('output', '') or str(remote_log_result)
    else:
        remote_log = str(remote_log_result) if remote_log_result else ''
    return {'user_email': user_email.strip(), 'remote_log': remote_log.strip()}

def get_thunderbird_folder_view__d9dcbaa59677e530494c93dfe8659800_qw35sft2_d8ca49d4(env, config: dict):
    """Read the folderTree mode from Thunderbird's xulstore.json."""
    xulstore_path = '/home/user/.thunderbird/t5q2a5hp.default-release/xulstore.json'
    file_bytes = env.controller.get_file(xulstore_path)
    if not file_bytes:
        return {'error': 'xulstore.json not found', 'mode': None}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        messenger = data.get('chrome://messenger/content/messenger.xhtml', {})
        folder_tree = messenger.get('folderTree', {})
        mode = folder_tree.get('mode', '')
        return {'mode': mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_thunderbird_multi_prefs__840f3ea7ded7f1b6695bc630a0a035b7_qw35sft2_ee0362cb(env, config: dict):
    """Read Thunderbird prefs.js and extract multiple preference values by key list."""
    import re
    pref_path = config.get('path', '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js')
    keys = config.get('keys', [])
    file_bytes = env.controller.get_file(pref_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    result = {}
    for key in keys:
        pattern = 'user_pref\\("' + re.escape(key) + '",\\s*(.*?)\\);'
        match = re.search(pattern, content)
        if match:
            raw_value = match.group(1).strip()
            if raw_value.startswith('"') and raw_value.endswith('"'):
                result[key] = raw_value[1:-1].replace('\\n', '\n').replace('\\"', '"')
            elif raw_value == 'true':
                result[key] = True
            elif raw_value == 'false':
                result[key] = False
            else:
                try:
                    result[key] = int(raw_value)
                except ValueError:
                    try:
                        result[key] = float(raw_value)
                    except ValueError:
                        result[key] = raw_value
        else:
            result[key] = None
    return result

def get_thunderbird_local_folders__1ea15ae7572bc843dca8cde6cbcd3a7c_qw35sft2_8c8fdc3d(env, config: dict):
    """
    Check which folders from the 'required_folders' list exist anywhere under
    the Thunderbird profile.  Returns a dict mapping folder name -> bool.
    """
    vm_ip = env.vm_ip
    port = env.server_port
    required = config.get('required_folders', ['COMPANY', 'UNIVERSITY', 'RESEARCH'])
    found = {}
    for folder in required:
        cmd = ['bash', '-c', f"find /home/user/.thunderbird -maxdepth 8 -name '{folder}' 2>/dev/null | head -1"]
        try:
            resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': cmd, 'shell': False}, timeout=15)
            if resp.status_code == 200:
                output = resp.json().get('output', '').strip()
                found[folder] = bool(output)
            else:
                found[folder] = False
        except Exception as e:
            logger_qw35sft2_7606d8.error('Error checking folder %s: %s', folder, e)
            found[folder] = False
    return found

def get_tb_forward_filter_state__f7fbe032caaa49b6e5ebdeafec739cf8_qw35sft2_916711f4(env, config: dict):
    """Get forward filter state from Thunderbird msgFilterRules.dat.

    Finds all msgFilterRules.dat files under .thunderbird, parses each filter
    block, and returns whether any filter has a forward action and the target
    email address.
    """
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not find_result or not find_result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_forward_filter': False, 'forward_to': None, 'filter_count': 0}
    filter_files = find_result['output'].strip().split('\n')
    has_forward = False
    forward_to = None
    filter_count = 0
    current_action = None
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
                current_action = None
            elif line.startswith('action='):
                current_action = line[7:].strip().strip('"').lower()
            elif line.startswith('actionValue=') and current_action is not None:
                val = line[12:].strip().strip('"')
                if 'forward' in current_action:
                    has_forward = True
                    forward_to = val
    return {'has_forward_filter': has_forward, 'forward_to': forward_to, 'filter_count': filter_count}

def get_thunderbird_smtp_no_incoming__f562ce0912b384f5748be41ba253880a_qw35sft2_0a0eb666(env, config: dict):
    """Get Thunderbird SMTP presence and absence of incoming mail accounts from prefs.js."""
    prefs_path = '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        bash_result = env.controller.run_bash_script('find /home/user/.thunderbird -name "prefs.js" -not -path "*/Crash*" 2>/dev/null | head -1', timeout=15)
        if bash_result and bash_result.get('stdout', '').strip():
            prefs_path = bash_result['stdout'].strip()
            file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    smtp_servers = re.findall('user_pref\\("mail\\.smtpservers",\\s*"([^"]+)"\\)', content)
    has_smtp = bool(smtp_servers and smtp_servers[0].strip())
    incoming_types = re.findall('user_pref\\("mail\\.server\\.\\w+\\.type",\\s*"([^"]+)"\\)', content)
    has_incoming = any((t in ('imap', 'pop3', 'nntp') for t in incoming_types))
    usernames = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.username",\\s*"([^"]+)"\\)', content)
    smtp_username = usernames[0] if usernames else ''
    return {'has_smtp': has_smtp, 'has_incoming': has_incoming, 'smtp_username': smtp_username}

def get_tb_acct_folder__2e9b3985bbecaa462dbe43eedf3d22fe_qw35sft2_479e0c2d(env, config: dict):
    """Check account removal from prefs.js and existence of a new local folder."""
    target_email = 'anonym-x2024@outlook.com'
    prefs_path = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    local_folders_dir = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders'
    r1 = env.controller.run_bash_script(f"grep -ic '{target_email}' '{prefs_path}' 2>/dev/null || echo 0", timeout=10)
    match_count = 0
    if r1:
        try:
            match_count = int(r1.get('output', '0').strip())
        except ValueError:
            match_count = -1
    account_removed = match_count == 0
    r2 = env.controller.run_bash_script(f"ls '{local_folders_dir}/' 2>/dev/null", timeout=10)
    folder_files = []
    if r2:
        folder_files = r2.get('output', '').splitlines()
    folder_name = 'Projects'
    folder_exists = any((f.strip() == folder_name or f.strip() == folder_name + '.msf' for f in folder_files))
    return {'account_removed': account_removed, 'folder_exists': folder_exists, 'folder_files': folder_files}

def get_thunderbird_draft_cc_attachment__ae64d5efe053bb1ed61ef261c2f7eddd_qw35sft2_b04c87df(env, config: dict):
    """Check Thunderbird Drafts MBOX for email with CC field and aws-bill.pdf attachment."""
    result_attach = env.controller.run_bash_script("find /home/user/.thunderbird -type f -name 'Drafts' 2>/dev/null | xargs grep -l 'aws-bill.pdf' 2>/dev/null | head -1", timeout=30)
    attach_output = ''
    if isinstance(result_attach, dict):
        attach_output = result_attach.get('output', '') or result_attach.get('stdout', '') or ''
    elif isinstance(result_attach, str):
        attach_output = result_attach
    draft_path = attach_output.strip()
    draft_with_attachment = bool(draft_path)
    draft_with_cc = False
    if draft_path:
        expected_cc = config.get('expected_cc', 'cfo@outlook.com')
        safe_cc = expected_cc.replace('"', '').replace("'", '')
        result_cc = env.controller.run_bash_script(f"grep -i 'Cc:.*{safe_cc}' '{draft_path}' 2>/dev/null | head -1", timeout=30)
        cc_output = ''
        if isinstance(result_cc, dict):
            cc_output = result_cc.get('output', '') or result_cc.get('stdout', '') or ''
        elif isinstance(result_cc, str):
            cc_output = result_cc
        draft_with_cc = bool(cc_output.strip())
    return {'draft_with_attachment': draft_with_attachment, 'draft_with_cc': draft_with_cc}

def get_thunderbird_imap_setup_state__7fc0ec566fbfa7a7ef758884ad90b95b_qw35sft2_02f51ec2(env, config: dict):
    """Get Account Setup dialog state, checking email and IMAP protocol presence."""
    try:
        tree = env.controller.get_accessibility_tree()
        if tree is None:
            return {'error': 'no_accessibility_tree', 'email_present': False, 'imap_present': False}
        tree_str = str(tree)
        email = config.get('expected_email', 'anonym-x2024@outlook.com')
        return {'email_present': email in tree_str, 'imap_present': 'IMAP' in tree_str, 'dialog_open': 'Account Setup' in tree_str}
    except Exception as e:
        return {'error': str(e), 'email_present': False, 'imap_present': False, 'dialog_open': False}

def get_tb_folder_two_filters__bc08452635875f98955858c473a86873_qw35sft2_574f5cb0(env, config: dict):
    """Check Thunderbird state: Promotions folder exists, discount filter exists, sale filter exists."""
    r = env.controller.run_bash_script('find /home/user/.thunderbird -maxdepth 1 -mindepth 1 -type d 2>/dev/null | grep -v "Crash Reports" | head -1', timeout=15)
    profile = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '').strip()
    if not profile:
        return {'promotions_exists': False, 'discount_filter': False, 'sale_filter': False}
    r2 = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Promotions" && echo "exists" || echo "missing"', timeout=15)
    promotions_out = r2.get('output', r2.get('stdout', '')) if isinstance(r2, dict) else r2 or ''
    promotions_exists = 'exists' in promotions_out
    r3 = env.controller.run_bash_script(f'find "{profile}" -name "msgFilterRules.dat" 2>/dev/null | xargs cat 2>/dev/null || echo ""', timeout=15)
    filter_content = (r3.get('output', r3.get('stdout', '')) if isinstance(r3, dict) else r3 or '').lower()
    discount_filter = 'discount' in filter_content
    sale_filter = 'sale' in filter_content
    return {'promotions_exists': promotions_exists, 'discount_filter': discount_filter, 'sale_filter': sale_filter}

def get_thunderbird_folder_view__0dc2b9b9129cecee63a42fd11a56a3f7_qw35sft2_342c6d1b(env, config: dict):
    """Read the folderTree mode from Thunderbird's xulstore.json."""
    xulstore_path = '/home/user/.thunderbird/t5q2a5hp.default-release/xulstore.json'
    file_bytes = env.controller.get_file(xulstore_path)
    if not file_bytes:
        return {'error': 'xulstore.json not found', 'mode': None}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        messenger = data.get('chrome://messenger/content/messenger.xhtml', {})
        folder_tree = messenger.get('folderTree', {})
        mode = folder_tree.get('mode', '')
        return {'mode': mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_thunderbird_nested_folders__1d59e1d1688abeb99fc9fce073e49b55_qw35sft2_f9dee36b(env, config: dict):
    """
    Check COMPANY, UNIVERSITY folders and an INBOX subfolder under each.
    In Thunderbird: subfolders of X live in X.sbd/ directory.
    Returns dict with boolean keys.
    """
    vm_ip = env.vm_ip
    port = env.server_port
    base_url = f'http://{vm_ip}:{port}/execute'
    try:
        resp = requests.post(base_url, json={'command': ['bash', '-c', "find /home/user/.thunderbird -maxdepth 8 -name 'COMPANY' 2>/dev/null | head -1"], 'shell': False}, timeout=15)
        company_out = resp.json().get('output', '').strip() if resp.status_code == 200 else ''
    except Exception as exc:
        logger_qw35sft2_c04067.error('Command failed: %s', exc)
        company_out = ''
    try:
        resp = requests.post(base_url, json={'command': ['bash', '-c', "find /home/user/.thunderbird -maxdepth 8 -name 'UNIVERSITY' 2>/dev/null | head -1"], 'shell': False}, timeout=15)
        university_out = resp.json().get('output', '').strip() if resp.status_code == 200 else ''
    except Exception as exc:
        logger_qw35sft2_c04067.error('Command failed: %s', exc)
        university_out = ''
    try:
        resp = requests.post(base_url, json={'command': ['bash', '-c', "find /home/user/.thunderbird -maxdepth 10 -name 'INBOX' -path '*/COMPANY.sbd/*' 2>/dev/null | head -1"], 'shell': False}, timeout=15)
        company_inbox_out = resp.json().get('output', '').strip() if resp.status_code == 200 else ''
    except Exception as exc:
        logger_qw35sft2_c04067.error('Command failed: %s', exc)
        company_inbox_out = ''
    try:
        resp = requests.post(base_url, json={'command': ['bash', '-c', "find /home/user/.thunderbird -maxdepth 10 -name 'INBOX' -path '*/UNIVERSITY.sbd/*' 2>/dev/null | head -1"], 'shell': False}, timeout=15)
        university_inbox_out = resp.json().get('output', '').strip() if resp.status_code == 200 else ''
    except Exception as exc:
        logger_qw35sft2_c04067.error('Command failed: %s', exc)
        university_inbox_out = ''
    return {'company': bool(company_out), 'university': bool(university_out), 'company_inbox': bool(company_inbox_out), 'university_inbox': bool(university_inbox_out)}

def get_tb_acct_smtp__93d97cb63d509700a5843e46760d10d0_qw35sft2_684ba59d(env, config: dict):
    """Check account removal and whether the associated SMTP server is also removed from prefs.js."""
    target_email = 'anonym-x2024@outlook.com'
    smtp_hostname = 'smtp.office365.com'
    prefs_path = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    r1 = env.controller.run_bash_script(f"grep -ic '{target_email}' '{prefs_path}' 2>/dev/null || echo 0", timeout=10)
    acct_count = 0
    if r1:
        try:
            acct_count = int(r1.get('output', '0').strip())
        except ValueError:
            acct_count = -1
    account_removed = acct_count == 0
    r2 = env.controller.run_bash_script(f"grep -ic '{smtp_hostname}' '{prefs_path}' 2>/dev/null || echo 0", timeout=10)
    smtp_count = 0
    if r2:
        try:
            smtp_count = int(r2.get('output', '0').strip())
        except ValueError:
            smtp_count = -1
    smtp_removed = smtp_count == 0
    return {'account_removed': account_removed, 'smtp_removed': smtp_removed, 'acct_match_count': acct_count, 'smtp_match_count': smtp_count}

def get_thunderbird_smtp_full__4c184097a04ef7d98ec36b30ff774f5c_qw35sft2_2427a295(env, config: dict):
    """Get all Thunderbird SMTP fields from prefs.js for full multi-step verification."""
    prefs_path = '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        bash_result = env.controller.run_bash_script('find /home/user/.thunderbird -name "prefs.js" -not -path "*/Crash*" 2>/dev/null | head -1', timeout=15)
        if bash_result and bash_result.get('stdout', '').strip():
            prefs_path = bash_result['stdout'].strip()
            file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    hostnames = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.hostname",\\s*"([^"]+)"\\)', content)
    usernames = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.username",\\s*"([^"]+)"\\)', content)
    ports = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.port",\\s*(\\d+)\\)', content)
    ssls = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.try_ssl",\\s*(\\d+)\\)', content)
    return {'hostname': hostnames[0] if hostnames else '', 'username': usernames[0] if usernames else '', 'port': int(ports[0]) if ports else 0, 'try_ssl': int(ssls[0]) if ssls else -1}

def get_thunderbird_filter_names__e046edd8a7ea67608e0814ad68f68ef0_qw35sft2_adcbd111(env, config: dict):
    """Find all message filter names from Thunderbird's msgFilterRules.dat."""
    result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not result or not result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'filter_names': []}
    filter_files = result['output'].strip().split('\n')
    filter_names = []
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                val = line[5:].strip()
                if val.startswith('"') and val.endswith('"'):
                    val = val[1:-1]
                filter_names.append(val)
    return {'filter_names': filter_names}

def get_thunderbird_compose_subject_attachment__59182577b48daf97ea5874cad0371f93_qw35sft2_1360ba35(env, config: dict):
    """Get subject and attachment status from open Thunderbird compose window."""
    tree = env.controller.get_accessibility_tree()
    tree_str = str(tree) if tree else ''
    expected_attachment = config.get('expected_attachment', 'aws-bill.pdf')
    expected_subject = config.get('expected_subject', 'May 2026 AWS Invoice')
    return {'has_attachment': expected_attachment in tree_str, 'has_subject': expected_subject in tree_str}

def get_tb_dual_filters__f4464353c01b135bac464c93385b0943_qw35sft2_42628857(env, config: dict):
    """Get dual filter state: forward filter and mark-as-read filter from Thunderbird.

    Reads all msgFilterRules.dat files and checks for the presence of:
    - A filter with a forward action (any actionValue)
    - A filter with a mark-as-read action
    Also captures the forward destination email if found.
    """
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not find_result or not find_result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_forward': False, 'has_mark_read': False, 'forward_to': None, 'filter_count': 0}
    filter_files = find_result['output'].strip().split('\n')
    filter_count = 0
    has_forward = False
    has_mark_read = False
    forward_to = None
    current_action = None
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
                current_action = None
            elif line.startswith('action='):
                val = line[7:].strip().strip('"').lower()
                current_action = val
                if 'mark' in val and 'read' in val:
                    has_mark_read = True
            elif line.startswith('actionValue=') and current_action:
                if 'forward' in current_action:
                    has_forward = True
                    forward_to = line[12:].strip().strip('"')
    return {'has_forward': has_forward, 'has_mark_read': has_mark_read, 'forward_to': forward_to, 'filter_count': filter_count}

def get_thunderbird_manual_cfg__598710d21fbdbbb5b9d697d62aaf5b30_qw35sft2_ae05246b(env, config: dict):
    """Get Manual Configuration view state, checking email and IMAP hostname."""
    try:
        tree = env.controller.get_accessibility_tree()
        if tree is None:
            return {'error': 'no_accessibility_tree', 'email_present': False, 'imap_host_present': False}
        tree_str = str(tree)
        email = config.get('expected_email', 'anonym-x2024@outlook.com')
        imap_host = config.get('expected_imap_host', 'outlook.office365.com')
        return {'email_present': email in tree_str, 'imap_host_present': imap_host in tree_str, 'manual_config_open': 'Manual configuration' in tree_str or 'INCOMING SERVER' in tree_str}
    except Exception as e:
        return {'error': str(e), 'email_present': False, 'imap_host_present': False, 'manual_config_open': False}

def get_tb_three_folders_one_filter__9865f0d026c5662c0f08b728c450437f_qw35sft2_3d152545(env, config: dict):
    """Check Thunderbird state: Promotions + Deals + Archive folders, discount filter."""
    r = env.controller.run_bash_script('find /home/user/.thunderbird -maxdepth 1 -mindepth 1 -type d 2>/dev/null | grep -v "Crash Reports" | head -1', timeout=15)
    profile = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '').strip()
    if not profile:
        return {'promotions_exists': False, 'deals_exists': False, 'archive_exists': False, 'discount_filter': False}
    r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Promotions" && echo "exists" || echo "missing"', timeout=15)
    promotions_out = r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or ''
    promotions_exists = 'exists' in promotions_out
    r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Deals" && echo "exists" || echo "missing"', timeout=15)
    deals_out = r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or ''
    deals_exists = 'exists' in deals_out
    r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Archive" && echo "exists" || echo "missing"', timeout=15)
    archive_out = r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or ''
    archive_exists = 'exists' in archive_out
    r = env.controller.run_bash_script(f'find "{profile}" -name "msgFilterRules.dat" 2>/dev/null | xargs cat 2>/dev/null || echo ""', timeout=15)
    filter_content = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '').lower()
    discount_filter = 'discount' in filter_content
    return {'promotions_exists': promotions_exists, 'deals_exists': deals_exists, 'archive_exists': archive_exists, 'discount_filter': discount_filter}

def get_thunderbird_multi_prefs__f88a527105940879c254b88e4c3fea68_qw35sft2_8c851feb(env, config: dict):
    """Read Thunderbird prefs.js and extract multiple preference values by key list."""
    import re
    pref_path = config.get('path', '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js')
    keys = config.get('keys', [])
    file_bytes = env.controller.get_file(pref_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    result = {}
    for key in keys:
        pattern = 'user_pref\\("' + re.escape(key) + '",\\s*(.*?)\\);'
        match = re.search(pattern, content)
        if match:
            raw_value = match.group(1).strip()
            if raw_value.startswith('"') and raw_value.endswith('"'):
                result[key] = raw_value[1:-1].replace('\\n', '\n').replace('\\"', '"')
            elif raw_value == 'true':
                result[key] = True
            elif raw_value == 'false':
                result[key] = False
            else:
                try:
                    result[key] = int(raw_value)
                except ValueError:
                    try:
                        result[key] = float(raw_value)
                    except ValueError:
                        result[key] = raw_value
        else:
            result[key] = None
    return result

def get_thunderbird_folder_view__3098a0959b134380c33dcaefd08a2ca7_qw35sft2_63010d62(env, config: dict):
    """Read the folderTree mode from Thunderbird's xulstore.json."""
    xulstore_path = '/home/user/.thunderbird/t5q2a5hp.default-release/xulstore.json'
    file_bytes = env.controller.get_file(xulstore_path)
    if not file_bytes:
        return {'error': 'xulstore.json not found', 'mode': None}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        messenger = data.get('chrome://messenger/content/messenger.xhtml', {})
        folder_tree = messenger.get('folderTree', {})
        mode = folder_tree.get('mode', '')
        return {'mode': mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_thunderbird_folder_state__50e2868517db7711c169fe6c468e7766_qw35sft2_9f561ac4(env, config: dict):
    """
    Check: COMPANY folder, UNIVERSITY folder, and a PROJECTS subfolder under COMPANY.
    In Thunderbird, sub-folders of COMPANY live in COMPANY.sbd/ directory.
    Returns dict with boolean keys.
    """
    try:
        res = env.controller.run_bash_script("find /home/user/.thunderbird -maxdepth 8 -name 'COMPANY' 2>/dev/null | head -1", timeout=15)
        company_out = res.get('output', '').strip() if isinstance(res, dict) else ''
    except Exception as e:
        logger_qw35sft2_848a1e.error('Command failed: %s', e)
        company_out = ''
    try:
        res = env.controller.run_bash_script("find /home/user/.thunderbird -maxdepth 8 -name 'UNIVERSITY' 2>/dev/null | head -1", timeout=15)
        university_out = res.get('output', '').strip() if isinstance(res, dict) else ''
    except Exception as e:
        logger_qw35sft2_848a1e.error('Command failed: %s', e)
        university_out = ''
    try:
        res = env.controller.run_bash_script("find /home/user/.thunderbird -maxdepth 10 -name 'PROJECTS' -path '*/COMPANY.sbd/*' 2>/dev/null | head -1", timeout=15)
        projects_out = res.get('output', '').strip() if isinstance(res, dict) else ''
    except Exception as e:
        logger_qw35sft2_848a1e.error('Command failed: %s', e)
        projects_out = ''
    return {'company': bool(company_out), 'university': bool(university_out), 'projects_in_company': bool(projects_out)}

def get_tb_acct_removal__3f1fce882b4fa7de95c0c939ed3f5b9d_qw35sft2_d7d94b11(env, config: dict):
    """Read prefs.js and check whether the target account email is absent."""
    target_email = 'anonym-x2024@outlook.com'
    prefs_path = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    result = env.controller.run_bash_script(f"grep -i '{target_email}' '{prefs_path}' 2>/dev/null | wc -l", timeout=10)
    if not result:
        return {'error': 'bash script failed', 'account_removed': False}
    output = result.get('output', '').strip()
    try:
        match_count = int(output)
    except ValueError:
        return {'error': f'unexpected output: {output}', 'account_removed': False}
    account_removed = match_count == 0
    return {'account_removed': account_removed, 'match_count': match_count}

def get_thunderbird_filter_match_all__f5a0e9992343a59dd61c9a37c78c3db0_qw35sft2_4b79d3bf(env, config: dict):
    """Check if any Thunderbird message filter uses 'Match all messages' (no conditions)."""
    result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not result or not result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_match_all_filter': False}
    filter_files = result['output'].strip().split('\n')
    has_match_all = False
    filter_count = 0
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
            elif line.startswith('condition='):
                val = line[10:].strip().strip('"').upper()
                if val == 'ALL':
                    has_match_all = True
    return {'has_match_all_filter': has_match_all, 'filter_count': filter_count}

def get_thunderbird_draft_attachment__ac90f87152777ec474b581760325d43f_qw35sft2_2768cbc8(env, config: dict):
    """Check Thunderbird Drafts MBOX for email containing aws-bill.pdf attachment."""
    result = env.controller.run_bash_script("find /home/user/.thunderbird -type f -name 'Drafts' 2>/dev/null | xargs grep -l 'aws-bill.pdf' 2>/dev/null | head -1", timeout=30)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '') or result.get('stdout', '') or ''
    elif isinstance(result, str):
        output = result
    draft_found = bool(output.strip())
    return {'draft_with_attachment': draft_found, 'draft_path': output.strip()}

def get_tb_incoming_forward__0ab002743507f6ef7d004b0ee40adbf4_qw35sft2_32a29ccc(env, config: dict):
    """Get forward filter with Getting New Mail trigger from Thunderbird.

    Reads msgFilterRules.dat and checks whether any filter has both:
    - type bit 0x1 set (Getting New Mail / Incoming trigger)
    - a forward action with a destination email
    """
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not find_result or not find_result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_incoming_forward': False, 'forward_to': None, 'filter_count': 0}
    filter_files = find_result['output'].strip().split('\n')
    filter_count = 0
    has_incoming_forward = False
    forward_to = None
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        current_type = None
        current_action = None
        current_action_val = None
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                if current_type is not None and current_type & 1 and current_action and ('forward' in current_action):
                    has_incoming_forward = True
                    if current_action_val:
                        forward_to = current_action_val
                filter_count += 1
                current_type = None
                current_action = None
                current_action_val = None
            elif line.startswith('type='):
                try:
                    current_type = int(line[5:].strip().strip('"'))
                except ValueError:
                    current_type = None
            elif line.startswith('action='):
                current_action = line[7:].strip().strip('"').lower()
            elif line.startswith('actionValue=') and current_action and ('forward' in current_action):
                current_action_val = line[12:].strip().strip('"')
        if current_type is not None and current_type & 1 and current_action and ('forward' in current_action):
            has_incoming_forward = True
            if current_action_val:
                forward_to = current_action_val
    return {'has_incoming_forward': has_incoming_forward, 'forward_to': forward_to, 'filter_count': filter_count}

def get_thunderbird_smtp_config__efa5fdf04c5026bfc6f8d6ffd452c5ec_qw35sft2_10c35c16(env, config: dict):
    """Get Thunderbird SMTP hostname and username from prefs.js."""
    prefs_path = '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        bash_result = env.controller.run_bash_script('find /home/user/.thunderbird -name "prefs.js" -not -path "*/Crash*" 2>/dev/null | head -1', timeout=15)
        if bash_result and bash_result.get('stdout', '').strip():
            prefs_path = bash_result['stdout'].strip()
            file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    hostnames = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.hostname",\\s*"([^"]+)"\\)', content)
    usernames = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.username",\\s*"([^"]+)"\\)', content)
    return {'hostname': hostnames[0] if hostnames else '', 'username': usernames[0] if usernames else ''}

def get_thunderbird_imap_port__c8735abb904c2556124a9fb1090745a7_qw35sft2_505ee3a8(env, config: dict):
    """Get Manual Configuration state, checking email, IMAP hostname, and port 993."""
    try:
        tree = env.controller.get_accessibility_tree()
        if tree is None:
            return {'error': 'no_accessibility_tree', 'email_present': False, 'imap_host_present': False, 'port_993_present': False}
        tree_str = str(tree)
        email = config.get('expected_email', 'anonym-x2024@outlook.com')
        imap_host = config.get('expected_imap_host', 'outlook.office365.com')
        return {'email_present': email in tree_str, 'imap_host_present': imap_host in tree_str, 'port_993_present': '993' in tree_str, 'manual_config_open': 'Manual configuration' in tree_str or 'INCOMING SERVER' in tree_str}
    except Exception as e:
        return {'error': str(e), 'email_present': False, 'imap_host_present': False, 'port_993_present': False, 'manual_config_open': False}

def get_tb_two_folders_two_filters__85999e7d7f538cb3e1b971a9ebe1ed0f_qw35sft2_a5976dbb(env, config: dict):
    """Check Thunderbird state: Promotions + Newsletter folders, discount + newsletter filters."""
    r = env.controller.run_bash_script('find /home/user/.thunderbird -maxdepth 1 -mindepth 1 -type d 2>/dev/null | grep -v "Crash Reports" | head -1', timeout=15)
    profile = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '').strip()
    if not profile:
        return {'promotions_exists': False, 'newsletter_exists': False, 'discount_filter': False, 'newsletter_filter': False}
    r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Promotions" && echo "exists" || echo "missing"', timeout=15)
    promotions_out = r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or ''
    promotions_exists = 'exists' in promotions_out
    r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Newsletter" && echo "exists" || echo "missing"', timeout=15)
    newsletter_out = r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or ''
    newsletter_exists = 'exists' in newsletter_out
    r = env.controller.run_bash_script(f'find "{profile}" -name "msgFilterRules.dat" 2>/dev/null | xargs cat 2>/dev/null || echo ""', timeout=15)
    filter_content = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '').lower()
    discount_filter = 'discount' in filter_content
    newsletter_filter = 'newsletter' in filter_content
    return {'promotions_exists': promotions_exists, 'newsletter_exists': newsletter_exists, 'discount_filter': discount_filter, 'newsletter_filter': newsletter_filter}

def get_thunderbird_local_folders__e760cb4cf85af456bb04229f8025e52c_qw35sft2_3c45b8b4(env, config: dict):
    """
    Check which folders from the 'required_folders' list exist anywhere under
    the Thunderbird profile.  Returns a dict mapping folder name -> bool.
    """
    vm_ip = env.vm_ip
    port = env.server_port
    required = config.get('required_folders', ['COMPANY', 'UNIVERSITY', 'WORK', 'PERSONAL'])
    found = {}
    for folder in required:
        cmd = ['bash', '-c', f"find /home/user/.thunderbird -maxdepth 8 -name '{folder}' 2>/dev/null | head -1"]
        try:
            resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': cmd, 'shell': False}, timeout=15)
            if resp.status_code == 200:
                output = resp.json().get('output', '').strip()
                found[folder] = bool(output)
            else:
                found[folder] = False
        except Exception as e:
            logger_qw35sft2_e8ef71.error('Error checking folder %s: %s', folder, e)
            found[folder] = False
    return found

def get_thunderbird_folder_view__582d00dc650f37f3e769d70699c5f1b6_qw35sft2_5b54ad33(env, config: dict):
    """Read the folderTree mode from Thunderbird's xulstore.json."""
    xulstore_path = '/home/user/.thunderbird/t5q2a5hp.default-release/xulstore.json'
    file_bytes = env.controller.get_file(xulstore_path)
    if not file_bytes:
        return {'error': 'xulstore.json not found', 'mode': None}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        messenger = data.get('chrome://messenger/content/messenger.xhtml', {})
        folder_tree = messenger.get('folderTree', {})
        mode = folder_tree.get('mode', '')
        return {'mode': mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_tb_acct_trash__8fc7e10e457380f7cb927cd320b9664e_qw35sft2_66bd48ec(env, config: dict):
    """Check account removal and whether Empty Trash on Exit is enabled in prefs.js."""
    target_email = 'anonym-x2024@outlook.com'
    prefs_path = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    r1 = env.controller.run_bash_script(f"grep -qi '{target_email}' '{prefs_path}' 2>/dev/null && echo 1 || echo 0", timeout=10)
    found_email = False
    if r1:
        found_email = r1.get('output', '0').strip() == '1'
    account_removed = not found_email
    r2 = env.controller.run_bash_script(f"grep -qi 'empty_trash_on_exit.*true' '{prefs_path}' 2>/dev/null && echo 1 || echo 0", timeout=10)
    empty_trash_enabled = False
    if r2:
        empty_trash_enabled = r2.get('output', '0').strip() == '1'
    return {'account_removed': account_removed, 'empty_trash_enabled': empty_trash_enabled}

def get_thunderbird_filter_and_pref__709b71706c15c6652aa94ca6ff4b6ae6_qw35sft2_34b2256c(env, config: dict):
    """
    Get both the applyIncomingFilters preference and filter count
    for multi-step evaluation (partial credit).
    """
    import re
    apply_filters = False
    prefs_result = env.controller.run_bash_script("cat /home/user/.thunderbird/t5q2a5hp.default-release/prefs.js 2>/dev/null | grep 'applyIncomingFilters'", timeout=10)
    if prefs_result and prefs_result.get('output', '').strip():
        line = prefs_result['output'].strip()
        m = re.search('applyIncomingFilters",\\s*(true|false)', line)
        if m and m.group(1) == 'true':
            apply_filters = True
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    filter_count = 0
    if find_result and find_result.get('output', '').strip():
        for fpath in find_result['output'].strip().split('\n'):
            fpath = fpath.strip()
            if not fpath:
                continue
            cat_result = env.controller.run_bash_script(f"grep -c '^name=' '{fpath}' 2>/dev/null || echo 0", timeout=10)
            if cat_result and cat_result.get('output', '').strip():
                try:
                    filter_count += int(cat_result['output'].strip())
                except ValueError:
                    pass
    return {'apply_incoming_filters': apply_filters, 'filter_count': filter_count, 'has_filter': filter_count > 0}

def get_tb_matchall_forward__bd4f2f9ff4f1fb7ec8f18ecd69248d51_qw35sft2_8473ec06(env, config: dict):
    """Get match-all condition and forward destination from Thunderbird msgFilterRules.dat.

    Checks whether any filter uses condition=ALL (match all messages) and
    whether any filter has a forward action with a destination email.
    """
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not find_result or not find_result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_match_all': False, 'forward_to': None, 'filter_count': 0}
    filter_files = find_result['output'].strip().split('\n')
    filter_count = 0
    has_match_all = False
    forward_to = None
    current_action = None
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
                current_action = None
            elif line.startswith('condition='):
                val = line[10:].strip().strip('"').upper()
                if val == 'ALL':
                    has_match_all = True
            elif line.startswith('action='):
                current_action = line[7:].strip().strip('"').lower()
            elif line.startswith('actionValue=') and current_action and ('forward' in current_action):
                forward_to = line[12:].strip().strip('"')
    return {'has_match_all': has_match_all, 'forward_to': forward_to, 'filter_count': filter_count}

def get_thunderbird_name_email__2d78e59033d5901e822a6624aeea3bc7_qw35sft2_60bf01ac(env, config: dict):
    """Get Account Setup dialog state, checking display name and email address."""
    try:
        tree = env.controller.get_accessibility_tree()
        if tree is None:
            return {'error': 'no_accessibility_tree', 'name_present': False, 'email_present': False}
        tree_str = str(tree)
        email = config.get('expected_email', 'anonym-x2024@outlook.com')
        name = config.get('expected_name', 'Anonym X')
        return {'name_present': name in tree_str, 'email_present': email in tree_str, 'dialog_open': 'Account Setup' in tree_str}
    except Exception as e:
        return {'error': str(e), 'name_present': False, 'email_present': False, 'dialog_open': False}

def get_thunderbird_compose_cc_attachment__433cf4a73d1a380efac6f447ba251498_qw35sft2_a2b5d268(env, config: dict):
    """Get CC field and attachment status from open Thunderbird compose window."""
    tree = env.controller.get_accessibility_tree()
    tree_str = str(tree) if tree else ''
    expected_attachment = config.get('expected_attachment', 'aws-bill.pdf')
    expected_cc = config.get('expected_cc', 'cfo@outlook.com')
    return {'has_attachment': expected_attachment in tree_str, 'has_cc': expected_cc in tree_str}

def get_thunderbird_smtp_description__8a8a5b00b6516b6e1171569c19648a1b_qw35sft2_cb59c0e2(env, config: dict):
    """Get Thunderbird SMTP server description from prefs.js."""
    prefs_path = '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        bash_result = env.controller.run_bash_script('find /home/user/.thunderbird -name "prefs.js" -not -path "*/Crash*" 2>/dev/null | head -1', timeout=15)
        if bash_result and bash_result.get('stdout', '').strip():
            prefs_path = bash_result['stdout'].strip()
            file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    descriptions = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.description",\\s*"([^"]+)"\\)', content)
    smtp_servers = re.findall('user_pref\\("mail\\.smtpservers",\\s*"([^"]+)"\\)', content)
    return {'description': descriptions[0] if descriptions else '', 'has_smtp': bool(smtp_servers and smtp_servers[0].strip())}

def get_tb_two_folders_and_filter__951712fea58d5a6835d2f3a50fd315b8_qw35sft2_cb18f390(env, config: dict):
    """Check Thunderbird state: Promotions + Deals folders exist, discount filter exists."""
    _r = env.controller.run_bash_script('find /home/user/.thunderbird -maxdepth 1 -mindepth 1 -type d 2>/dev/null | grep -v "Crash Reports" | head -1', timeout=15)
    profile = (_r.get('output', _r.get('stdout', '')) if isinstance(_r, dict) else _r or '').strip()
    if not profile:
        return {'promotions_exists': False, 'deals_exists': False, 'discount_filter': False}
    _r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Promotions" && echo "exists" || echo "missing"', timeout=15)
    promotions_out = _r.get('output', _r.get('stdout', '')) if isinstance(_r, dict) else _r or ''
    promotions_exists = 'exists' in promotions_out
    _r = env.controller.run_bash_script(f'test -e "{profile}/Mail/Local Folders/Deals" && echo "exists" || echo "missing"', timeout=15)
    deals_out = _r.get('output', _r.get('stdout', '')) if isinstance(_r, dict) else _r or ''
    deals_exists = 'exists' in deals_out
    _r = env.controller.run_bash_script(f'find "{profile}" -name "msgFilterRules.dat" 2>/dev/null | xargs cat 2>/dev/null || echo ""', timeout=15)
    filter_content = _r.get('output', _r.get('stdout', '')) if isinstance(_r, dict) else _r or ''
    discount_filter = 'discount' in filter_content.lower()
    return {'promotions_exists': promotions_exists, 'deals_exists': deals_exists, 'discount_filter': discount_filter}

def get_tb_theme_and_folder__bdfb5942d44e1763e7f54cb6f6b53816_qw35sft2_217de1fc(env, config: dict):
    """Get active Thunderbird theme ID and check if Night Notes local folder exists."""
    import re
    import json
    profile_prefs = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    folder_msf = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Night Notes.msf'
    theme_id = None
    try:
        prefs_bytes = env.controller.get_file(profile_prefs)
        if prefs_bytes:
            prefs_text = prefs_bytes.decode('utf-8', errors='replace')
            pref_re = re.compile('^user_pref\\("extensions\\.activeThemeID",\\s*(.+)\\);$', re.MULTILINE)
            m = pref_re.search(prefs_text)
            if m:
                try:
                    theme_id = json.loads(m.group(1))
                except Exception:
                    theme_id = m.group(1).strip('"')
    except Exception:
        pass
    folder_exists = False
    try:
        msf_bytes = env.controller.get_file(folder_msf)
        folder_exists = msf_bytes is not None
    except Exception:
        folder_exists = False
    return {'active_theme_id': theme_id, 'night_notes_folder_exists': folder_exists}

def get_thunderbird_folder_view__9ec4687f5770fc31e26559c7787fe14b_qw35sft2_8b914604(env, config: dict):
    """Read the folderTree mode from Thunderbird's xulstore.json."""
    xulstore_path = '/home/user/.thunderbird/t5q2a5hp.default-release/xulstore.json'
    file_bytes = env.controller.get_file(xulstore_path)
    if not file_bytes:
        return {'error': 'xulstore.json not found', 'mode': None}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        messenger = data.get('chrome://messenger/content/messenger.xhtml', {})
        folder_tree = messenger.get('folderTree', {})
        mode = folder_tree.get('mode', '')
        return {'mode': mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_thunderbird_filter_incoming__4d2f365259c66d4c891bf0c4c8a5a5f0_qw35sft2_5bcea5d8(env, config: dict):
    """Check if any Thunderbird message filter has 'Getting New Mail' (Incoming) trigger enabled."""
    result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not result or not result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'has_incoming_filter': False, 'filter_count': 0}
    filter_files = result['output'].strip().split('\n')
    has_incoming = False
    filter_count = 0
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        lines = content.splitlines()
        current_enabled = True
        current_type = None
        for line in lines:
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
                current_enabled = True
                current_type = None
            elif line.startswith('enabled='):
                val = line[8:].strip().strip('"')
                current_enabled = val.lower() == 'yes'
            elif line.startswith('type='):
                try:
                    current_type = int(line[5:].strip().strip('"'))
                except ValueError:
                    current_type = None
                if current_enabled and current_type is not None and current_type & 1:
                    has_incoming = True
    return {'has_incoming_filter': has_incoming, 'filter_count': filter_count}

def get_tb_acct_filter__ddc543b6ee27062796019a5514b0e693_qw35sft2_c9574301(env, config: dict):
    """Check account removal and whether a 'newsletter' message filter exists in Local Folders."""
    target_email = 'anonym-x2024@outlook.com'
    prefs_path = '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js'
    r1 = env.controller.run_bash_script(f"grep -ic '{target_email}' '{prefs_path}' 2>/dev/null || echo 0", timeout=10)
    match_count = 0
    if r1:
        try:
            match_count = int(r1.get('output', '0').strip())
        except ValueError:
            match_count = -1
    account_removed = match_count == 0
    r2 = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | xargs grep -il 'newsletter' 2>/dev/null | wc -l", timeout=10)
    newsletter_filter_count = 0
    if r2:
        try:
            newsletter_filter_count = int(r2.get('output', '0').strip())
        except ValueError:
            newsletter_filter_count = 0
    newsletter_filter_exists = newsletter_filter_count > 0
    return {'account_removed': account_removed, 'newsletter_filter_exists': newsletter_filter_exists}

def get_tb_named_forward_filter__e2cb27642e2d718992a0ab3ec3b240f1_qw35sft2_56f302c1(env, config: dict):
    """Get filter names and forward destination from Thunderbird msgFilterRules.dat.

    Returns the list of filter names found and the last forwarding destination
    email encountered, to allow checking both filter name and target.
    """
    find_result = env.controller.run_bash_script("find /home/user/.thunderbird -name 'msgFilterRules.dat' 2>/dev/null | head -5", timeout=10)
    if not find_result or not find_result.get('output', '').strip():
        return {'error': 'msgFilterRules.dat not found', 'filter_count': 0, 'filter_names': [], 'forward_to': None}
    filter_files = find_result['output'].strip().split('\n')
    filter_count = 0
    filter_names = []
    forward_to = None
    current_action = None
    for fpath in filter_files:
        fpath = fpath.strip()
        if not fpath:
            continue
        read_result = env.controller.run_bash_script(f"cat '{fpath}' 2>/dev/null", timeout=10)
        if not read_result or not read_result.get('output'):
            continue
        content = read_result['output']
        for line in content.splitlines():
            line = line.strip()
            if line.startswith('name='):
                filter_count += 1
                filter_names.append(line[5:].strip().strip('"'))
                current_action = None
            elif line.startswith('action='):
                current_action = line[7:].strip().strip('"').lower()
            elif line.startswith('actionValue=') and current_action and ('forward' in current_action):
                forward_to = line[12:].strip().strip('"')
    return {'filter_count': filter_count, 'filter_names': filter_names, 'forward_to': forward_to}

def get_thunderbird_compose_bcc_attachment__d67afa1e0d57a9bbcae95982d894cc83_qw35sft2_644cc038(env, config: dict):
    """Get BCC field and attachment status from open Thunderbird compose window."""
    tree = env.controller.get_accessibility_tree()
    tree_str = str(tree) if tree else ''
    expected_attachment = config.get('expected_attachment', 'aws-bill.pdf')
    expected_bcc = config.get('expected_bcc', 'finance@outlook.com')
    return {'has_attachment': expected_attachment in tree_str, 'has_bcc': expected_bcc in tree_str}

def get_thunderbird_smtp_security__9cc78735939880d6658a6582ff470dcd_qw35sft2_e2f5e05f(env, config: dict):
    """Get Thunderbird SMTP port and connection security (try_ssl) from prefs.js."""
    prefs_path = '/home/user/.thunderbird/6ex3j72p.default-release/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        bash_result = env.controller.run_bash_script('find /home/user/.thunderbird -name "prefs.js" -not -path "*/Crash*" 2>/dev/null | head -1', timeout=15)
        if bash_result and bash_result.get('stdout', '').strip():
            prefs_path = bash_result['stdout'].strip()
            file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        return {'error': 'prefs.js not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    ports = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.port",\\s*(\\d+)\\)', content)
    ssls = re.findall('user_pref\\("mail\\.smtpserver\\.\\w+\\.try_ssl",\\s*(\\d+)\\)', content)
    return {'port': int(ports[0]) if ports else 0, 'try_ssl': int(ssls[0]) if ssls else -1}

def get_tb_two_folders_two_filters_v2__627b4b586c777792c6f3d931dd087820_qw35sft2_cf0ddccb(env, config: dict):
    """Check Thunderbird state: Promotions + Coupons folders, discount + coupon filters."""
    run_cmd = lambda cmd: (lambda r: r.get('output', r.get('stdout', '')) if isinstance(r, dict) else r or '')(env.controller.run_bash_script(cmd, timeout=15))
    profile = run_cmd('find /home/user/.thunderbird -maxdepth 1 -mindepth 1 -type d 2>/dev/null | grep -v "Crash Reports" | head -1').strip()
    if not profile:
        return {'promotions_exists': False, 'coupons_exists': False, 'discount_filter': False, 'coupon_filter': False}
    promotions_out = run_cmd(f'test -e "{profile}/Mail/Local Folders/Promotions" && echo "exists" || echo "missing"')
    promotions_exists = 'exists' in promotions_out
    coupons_out = run_cmd(f'test -e "{profile}/Mail/Local Folders/Coupons" && echo "exists" || echo "missing"')
    coupons_exists = 'exists' in coupons_out
    filter_content = run_cmd(f'find "{profile}" -name "msgFilterRules.dat" 2>/dev/null | xargs cat 2>/dev/null || echo ""').lower()
    discount_filter = 'discount' in filter_content
    coupon_filter = 'coupon' in filter_content
    return {'promotions_exists': promotions_exists, 'coupons_exists': coupons_exists, 'discount_filter': discount_filter, 'coupon_filter': coupon_filter}

def get_ext_and_minimap__8a991aa9f9913ccca00f4b1c76aac764_qw35sft2_88c50775(env, config: dict):
    """Get extension list and editor.minimap.enabled from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_52d4f8.warning('Failed to get extension list: %s', e)
    minimap_enabled = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            minimap_enabled = settings.get('editor.minimap.enabled')
    except Exception as e:
        logger_qw35sft2_52d4f8.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'minimap_enabled': minimap_enabled}
