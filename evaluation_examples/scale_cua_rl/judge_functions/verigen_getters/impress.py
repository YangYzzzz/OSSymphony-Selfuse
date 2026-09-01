"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_impress_title_color__9e6f539665d909e9d4b70212d289e0ee', 'get_pptx_slide_title__ebdab037a8eab6c57e56f5796ba7c0ab', 'get_pptx_multi_slide_images__e6e9ce905cbd4fd74addd1fc2cd640de', 'get_impress_save_settings__a73132a8bbdfe6ee96f72881e60f8547', 'get_impress_title_bold__675486c95c0c9441584aa662f56cd677', 'get_slide_image_count__07113500ab2edd0ce92ed5a5e9e8b1aa', 'get_pptx_table_rows__d7904b6a5ae76e528f93447180e8419d', 'get_pptx_shape_dimensions__30ab1f17976bddf578367b54740da469', 'get_impress_slide_table_info__dabff5cb0982227a6c01462d5c667c85', 'get_pptx_slide_count__3e661325ce185501287a57bdd6c09b52', 'get_pptx_slide_title__95b2513b7f5b16d49b143853ab08b964', 'get_impress_title_underline__b1d7e7e1d48378504a47d3ccfb3cf1be', 'get_pptx_slide_title_and_bg__93a38408d9330c4f554c19c28759d1cc', 'get_impress_slide_title__497f89854667fa402a3b0a0e21fc061d', 'get_pptx_text_alignment__d36eb73e90de91fc69ee1df98b9d5058', 'get_pptx_table_exists__ffa5734cd742b977903cbddf387ddf6c', 'get_pptx_slide_title_text__661ea64646d6a38833b5fdf5450470f8', 'get_impress_middle_text_font_size__e0ab347b041b8b1f40c9ca0e5bc35d39', 'get_impress_slide_bg_color__147ec12b6038f2d22f20700395a41639', 'get_pptx_slide_audio__1b19ca869f4a35a02b8d110c7e17b7e9', 'get_pptx_slide_title__36872752e18a684ac0b72ebb0b9e2c33', 'get_impress_table_bg__9da3068d76a130c43b0820c0580d6baa', 'get_impress_content_bold__ca58c5a207cc4d7f5ab415f13aba66b0', 'get_pptx_table_row__d2dd546208d48aeec32a49334e180bbc', 'get_impress_image_position__b7ad7d38cb36c439bba7c5f1f082cb99', 'get_impress_bg_fontsize__c22d8b15c916d1ef420eb42daeb27bf0', 'get_pptx_slide1_title_format__7d9f2200326202f628d33a5691960f9c', 'get_pptx_text_bold__c4941aca377cb2eab0e421e15f393854', 'get_slide_orientation__662bf4c59a75818f259c7e0aef709493', 'get_pptx_title_font_size__a1c0faa26bf29b13cb4b3dfa5122d14a', 'get_impress_title_style_bg__d816da51e2714a5b86252b84b6e77125', 'get_slide1_image_position__ebf21364e626ee60bab03f3d9d12deff', 'get_impress_subtitle_font_color__d0c681af3bd2af47798202cb5ebe1988', 'get_slide_text_font_color__3892513c1cab5055660dfeaa188121fb', 'get_pptx_slide_title__7b96cb41298a29a8e5448750424ab134', 'get_impress_strikethrough__a345b8271c8198e8c95ec913d78c781f', 'get_pptx_pic_width_and_fonts__1a124508c5c52deecd8c9cfabb1cdd16', 'get_impress_slide_info__b84f4f9ed96e0e954dfada4f0a280734', 'get_pptx_title_font_color__3536f82770e6c6d12bb7de4bc43120d5', 'get_pptx_image_props__32c4d5ced54cde2e00ee07b6ba3973e4', 'get_impress_bg_and_subtitle__024a68353adff6dd2d6a01cc842a46ed', 'get_slide_title_font_color__66374a642830d11922d74321b71d53a4', 'get_pptx_title_format__322bfa5ed10c48cee5b1ded75b68bcd6', 'get_pptx_shape_dimensions__37468b49cfdc53eaca3dfbbee86985a3', 'get_pptx_title_text__927f9acbea3f73fbaa9e281542feb428', 'get_impress_content_align_title_color__63833ad040c5e27e84bd66675fcf45fe', 'get_slide_count__55d37620fcf7998461b41e2b1579f16b', 'get_impress_title_bold__6e02436ff918c1445d486be48242a3ca', 'get_impress_slide_info__1abdd5c6fe686f66f2606ee0a55bf0d1', 'get_pptx_textbox_font_sizes__a0ffb78fbe8b896461c46fcb419a2efb', 'get_impress_all_title_fonts__9c7bd04930deae5cddd57fa6475996f4', 'get_slide_count_and_layout__c9d9670d65b42979a3ba7969e086139a', 'get_impress_text_and_bg__6d65cdf6ba3aad0bcf7dd7e1b9c52c8f', 'get_impress_title_fontsize__f67d7e69d6a85882636e0a41dabd5c89', 'get_pptx_text_underline__1563ed9d69890bfb1ca48bd048aeee4a', 'get_pptx_table_bold__1f6e651edc27326d3dac73dc9e29333b', 'get_pptx_textbox_font_sizes__8e211b70307c1527c1dc374c2a15baa9', 'get_pptx_text_alignment__b3837f805dc03cde831accc6a7063290', 'get_pptx_orientation__2f288a1714ed5a4b53fed9b568a61679', 'get_pptx_bold_check__c8d3ea84bab1df08b4dcc2adb9177ac5', 'get_pptx_textbox_font_sizes__d842f0d6650bc2853fbd82b34d9bb662', 'get_pptx_slide_info__e9203909bdddcf8283496218733a0b9b', 'get_slide_notes__44fd5b3a7c457950926c4afa82f02524', 'get_pptx_slide_title__07d639643e02063bb060c675d5eb936b', 'get_pptx_slide2_welcome_format__203dc282bb468063e328a7afef0f556d', 'get_pptx_slide_text__64fdb0e52d9e147a698cbe6791191c81', 'get_impress_bg_font_color__d8eda1446480065d3fccc5f3a4fe966e', 'get_impress_top_text_color__42394c3fddd5c54d80cb5da5e81a07d8', 'get_impress_table_with_content__608b93868302b656e50a31095c9d18e9', 'get_slide_orientation__060b6ae42d4d041f46057b0dff06cb53', 'get_impress_notes_bg__7bae80483f423bfa3b3bbb7f9e9eafdb', 'get_slide_count__1573725569136f71d4ad79ac5ffbf35a', 'get_slide_bg_color__1dbf2adcf932ab2d0e517ac3d78f9d08', 'get_impress_slide_info__88930a069947104f21bb2b61fdebcc2c', 'get_impress_title_color__6b28d17bfe07bb684917211f85bf497c', 'get_slide_image_info__e8973b7aba53f626fa037406b1da4c8e', 'get_impress_title_bg__c2ac2f94e6c641c5a5c2e733ff74ade1', 'get_pptx_slide_paragraphs__10c4d6d7b532dd358d2936a84f08db4a', 'get_snapshot_and_slide_bg__da1a5547fec970e705ee94d1f2bd1e91', 'get_pptx_slide_dimensions__54270a3c8d106e03ce2d0c72089a3417', 'get_impress_bottom_text_bold__06366a1cfa1039e4258f94612f6055b9', 'get_pptx_pic_height_and_fonts__4eb9ebb1f165db3d3eb60748b0824590', 'get_pptx_slide_text__38a0fe7ffb236efbdbbcf3cc7460e807', 'get_slide_audio_info__5688a332db768af3543ab15e11555c76', 'get_impress_title_font_props__8c5c9bdd88c0ba0d1f43df42a252cfae', 'get_pptx_title_bold__7728836b00929340fdef9a2a6602ef09', 'get_impress_slide2_state__33bd53bce176d3722b0b7a880cb5489e', 'get_slide_text_font_color__175bcb8418fdbfeda65db92690d9b128', 'get_slide_text_bold__db0f57f6a04c8947ec0364b499605e29', 'get_impress_title_format__110ae328b957703cc4b8db2a2b04df5c', 'get_impress_strikethrough__bbd5669ff0c1c503844620466c04ffda', 'get_slide_notes__4bc42c00ac052a9214daa0b289a72718', 'get_pptx_last_slide_info__964158cb23267c07918d09a84fcd28cc', 'get_slide_orientation__326a96d7a18db487b955505da7fa674b', 'get_pptx_italic_check__979792a60bf386d8e1c5a4d702d19829', 'get_pptx_text_alignment__a27e67f69ab078486669e5f21f103d16', 'get_pptx_slide_hidden_status__266b91e28f8289ba603f3990e8aa8eb6', 'get_impress_slide_picture_count__48004e0f145da635f1d1ef72ec0c75ce', 'get_pptx_slide_subtitle__b0063e78116208a0699a2e91287a92ca', 'get_pptx_pic_height_and_fonts__5b53387f975ab04f84e7061a6b89939a', 'get_impress_title_size_content_bold_bg__80858354e7972b18d1e56bf385980263', 'get_impress_title_font_color__dcc119efe7774eb3168a0af998c68532', 'get_impress_slide_title__94c866527f562343eb1ed6d961346913', 'get_pptx_slide_count_and_titles__b3862aa49d8968cf4ed21032ba8b229f', 'get_impress_slide_count_and_titles__3a9a25804cb58719a25a23e52e7e124a', 'get_pptx_picture_position__77f9e57f0af70c9666efd508929f72c3', 'get_pptx_slide_title__6c6be158040f1b190133813bff96d3a4', 'get_pptx_slide_count__f16a19b8aecc56a3153baf7aa9668d95', 'get_impress_slide_title_text__1b7b08a28b40e394d9044a5b56ecfe62', 'get_pptx_image_props__19c7a30018144cafc9f4ebd4b22910ef', 'get_slide1_badge_exists__6b4bf280e2d30f49e851d7a1a643ef02', 'get_pptx_text_italic__39e2309fd148a4ecc88950c3fc66ef3b', 'get_impress_slide_title_font__05b860587b204cf23347a8d6b2a21320', 'get_slide1_image_width__d195f7f7ed997b32944186b135a8c2d3', 'get_slide_title_font__7f1ba412564f3dc4322d685ce4a3274f', 'get_impress_title_font_size__3af9d622838a82de0731abee0010dddc', 'get_pptx_slide5_title_format__e80588656c5c992160d4d54a69813615', 'get_pptx_slide_title_with_color__fd2e510c274f3f89d65383d731b18828', 'get_impress_italic_subtitle__910b283898dbff47da6b9e5280c358bb', 'get_pptx_slide_bg_color__6efa38f597117d6ec1022844076cbdb0', 'get_slide_has_bg_image__54b7105f00ae5aefd1a4f4e9a1686c3b', 'get_pptx_slide_title__24e166a4a3628d77570c2919341fa3d0', 'get_impress_strikethrough__941b469f6d56a3f2848b88f32a6979cb', 'get_pptx_transitions__a8621b66c753635f21d8c612a865d275', 'get_pptx_text_underline__935b1aaa61bb5f50d9e66c69d9b5829b', 'get_pptx_shape_dimensions__dc32e28ff1403085a7367144a8c3b449', 'get_pptx_slide_title__f01c4fcc320219e33dd1b3dfc752a5b6', 'get_pptx_orientation__f921e5b58a724c7baa4b59415c1c4019', 'get_impress_title_and_bg__ac120f719f06d988c1c09c7022ee4b50', 'get_pptx_slide_count_and_title__e6d2505b7a2f1323c83204293d9b2ff3', 'get_impress_strikethrough__0be9a8c316c2b4d9a9ba4471962b751f', 'get_pptx_all_slides_bg_color__88a6aaaf9999692a07033829e3809d82', 'get_pptx_table_state__e5c61d88a7cca7248552b922fe7c3602_qw35sft2_3635603e', 'get_pptx_slide_count__d96a0661e83ae9239b523060c49e36f2_qw35sft2_982cb0b0', 'get_pptx_table_row0__27882b2f85f9d8345d5e8153a8f07379_qw35sft2_49026ab3', 'get_pptx_summary_notes__f1eef050d638a5408c254c4e620a5302_qw35sft2_7b56b6c6', 'get_pptx_first_slide_text__b9dba6723fe45eeec63ed20d4f46d22d_qw35sft2_c385ede6', 'get_impress_slide2_title_props__dea2c3a0c75be4ec786c2a7bbd9ea59f_qw35sft2_861497b7', 'get_pptx_slide_transitions__f348aaf139c455284fcbda7cfe577da5_qw35sft2_4f3addf3', 'get_impress_file_and_transition__5e93f9ddd16cd129b875ad7386e8a5e5_qw35sft2_a0fba5ee', 'get_pptx_slide_image_info__7a22b4a6c189fa88dd4ad54ac62a04cd_qw35sft2_2b48f3ad', 'get_impress_presenter_console_state__c6ff7e494d134a42d6e5c420e00bf4fd_qw35sft2_0e3dd544', 'get_pptx_bg_fill_type__7f279c8dbe68c5780c5b5f60310b951b_qw35sft2_67387c51', 'get_pptx_transitions__85331dc0cf4d821fa2bad3e27657d683_qw35sft2_0149f70f', 'get_pptx_slide_bg_color__1a5b2a2b44ee739e3486ddbc20ac1c87_qw35sft2_88959acb', 'get_pptx_slide_bg_fill__30ec500dc9631258b0ce307c50b98145_qw35sft2_26833ac2', 'get_impress_slide4_contact_shapes__5eae632a29df23abd3f96eee0032e6dc_qw35sft2_c546a846', 'get_pptx_stretch_and_shapes__bd30ea01dba8f440569b065c94e4f32f_qw35sft2_7b4096a9', 'get_impress_underline_state__fa317ec6034191f7b02a906f78ff4b09_qw35sft2_128f5711', 'get_impress_audio_slide1_trans_slide3__1b01b098ee4b6dfab8bc1b2d24723a1f_qw35sft2_18b7000f', 'get_impress_slide_info__da57195bec8b00f5d49d80dfe28b0049_qw35sft2_ccd6c0f7', 'get_impress_bullet_newpara__12c621423b2418d27e86fd746731530a_qw35sft2_457c566e', 'get_pptx_text_alignments__6253f908caa7dec8dc9ccbad6c2f7332_qw35sft2_c050a0a8', 'get_pptx_slide_font_size__4b92c85bf54981a99e31cdf8daa555a4_qw35sft2_86e22296', 'get_pptx_slide3_text_colors__9e2b961f79a8b89da2fe1a7f40c57f65_qw35sft2_4dc70f53', 'get_pptx_slide3_para_texts__c0cb167c6059fd8f9f81e328b9599144_qw35sft2_81111149', 'get_impress_multi_title_color__14718ed4f8dc81749072a781b972112d_qw35sft2_09d2b40d', 'get_slide2_notes_state__dd477a39d59856cb4084c0b8f3f8bce8_qw35sft2_a6f6bdb8', 'get_impress_slide1_title__992b8b598a921e5736a852e261be2237_qw35sft2_7462f87d', 'get_pptx_image_quadrant__f123db274ed42f1e244c71177f056c09_qw35sft2_5f5fc1c9', 'get_pptx_slide2_title_state__1225c3617161e6b819a6b6275a5c9749_qw35sft2_b7f0038d', 'get_impress_slide3_table_pos__db4d3e235d7b0e304c073a25921cbfa2_qw35sft2_c36fb8ca', 'get_pptx_title_underline__42b6fd5c06cd718a123cb3a7892d354e_qw35sft2_fe20d734', 'get_impress_slide14_font_sizes__ec3f2f9f447d657b0668843895ae7804_qw35sft2_1714f192', 'get_impress_notes_bg_title__cb759dd1cf89c2aec2802f8543d62400_qw35sft2_73e00eef', 'get_pptx_text_color__39fe73dd35955e4e54b2d721092898e1_qw35sft2_47e52091', 'get_pptx_single_text_color__c830024b32b30956e96002af134f6f4a_qw35sft2_25f7b637', 'get_slide2_bg_subtitle_color__2d907c4b088ee7e1bdf99fbd932517ff_qw35sft2_3d0a0bfd', 'get_pptx_title_subtitle_font__761afd0460bdb180560f4c7116933937_qw35sft2_0b2c4583', 'get_impress_pptx_props__ed92c6da7770c48f6ebca23c6070cbf5_qw35sft2_da03d5c6', 'get_pptx_content_and_table__7f43e77e7e3ce83ed182df2ac03d3d3a_qw35sft2_1259b210', 'get_picture_size_slide6__f55d52551229ea682c89b1a15474ccff_qw35sft2_cf035c45', 'get_impress_panel_slide_count__ddf9714ddb153a97bbdfb4350e733817_qw35sft2_1a390bcc', 'get_pptx_text_font_size__f00d8a2d9828be4b94677a9606cd7f47_qw35sft2_37d00a00', 'get_pptx_orientation__2491631e79810e0b815ac32866fd3274_qw35sft2_0f4fbcfa', 'get_pptx_orientation_state__f5c58fc75d059d716d7a9ec9e8e3967d_qw35sft2_e2628c21', 'get_pptx_table_state__99738e8866681b9fdc97e255c87e124d_qw35sft2_2395336c', 'get_impress_slide2_title_top__ff47db79abc78d9904bd7e17dbb1e3d3_qw35sft2_13ea5fa8', 'get_pptx_notes_nonempty__d897409e855e2cbbb791ebde03eeef8a_qw35sft2_37fdad70', 'get_pptx_slide_title_text__c928e38a883d919c295367c551657155_qw35sft2_7b8b2632', 'get_pptx_slide_image_and_transition__4195ceef20cb2389eadfdf0daf18e8e7_qw35sft2_77e2d736', 'get_pptx_slide_bg_color__b309be1bc27dcbe28bf5749a28cc9b6f_qw35sft2_f940a712', 'get_pptx_slide_subtitle_text__4d840f044cdf667c54faf806c6af73e1_qw35sft2_3db23da2', 'get_pptx_last_two_duplicated__ae132e5d53d956fce4051acdf425fcba_qw35sft2_499adc83', 'get_pptx_stretch_and_transition__676dd2418e7e3a84c9083aa5ad9f76e2_qw35sft2_861ec0f5', 'get_impress_slide_info__881be08815c16b0f7ba88245e2d69e10_qw35sft2_207240ca', 'get_impress_audio_and_transition__1b19ca869f4a35a02b8d110c7e17b7e9_qw35sft2_6f055a7e', 'get_pptx_slide_font_italic__5e9f703dc627e70a3b4d09452ec3d3ab_qw35sft2_b38d45fc', 'get_pptx_slides35_colors__0fbb4b3515e1009c69bc08919cd2b491_qw35sft2_1cbad4a0', 'get_pptx_slide3_subpoint_level__3c50fb39e50d5c4bd155f8f4b7e911d1_qw35sft2_170ff70f', 'get_impress_file_slide_count__e4d3506bfaa42b2028f1a93213c30bc6_qw35sft2_f8e2a094', 'get_impress_bullet_underline__396e098239a8be5c53a96d982334cb7b_qw35sft2_75a01ec0', 'get_impress_options_combo__43da862b24de1bb5002a875bb93d754e_qw35sft2_7a41a212', 'get_pptx_slide_orientation__1830f729a2ae0637260a8bca58b0c8cd_qw35sft2_29845848', 'get_impress_img_pos__edae955a62ba6a9d5bff7ce774e3ee10_qw35sft2_c281af6e', 'get_impress_title_color_underline__84ce98fac5a701a1ff43f1337a4f838c_qw35sft2_8e6a55dc', 'get_pptx_slide2_title_bold__ba024eb57784584ce2b08e31db5871ac_qw35sft2_538afce5', 'get_pptx_image_position__2e5aa88aaf11fa4bda2a55acfcab9dc5_qw35sft2_9a585c0e', 'get_slide2_title_state__78dfac329c5cd8b1c5efec5b79600fdd_qw35sft2_c4a7bb63', 'get_impress_slide3_title_bold__7dcab435f4cba9d9e23699eed9a4d408_qw35sft2_11c65a9e', 'get_pptx_text_alignments__14ff9da6ca84bfee11448e7c7f8efafc_qw35sft2_5c719df5', 'get_impress_notes_bg_transition__5cb083f40f95f75316141f060ee8a1cc_qw35sft2_1fc3c009', 'get_impress_slide14_font_bold__a28710b7c875e09d23efa44313d2c747_qw35sft2_4643d71f', 'get_impress_slide3_table_and_title__80d9e8d7efd024cbc51991913269b5fb_qw35sft2_53d7293b', 'get_impress_slide_title_full__9620e3371e8842ebb86c9c7ef6cb4f00_qw35sft2_c6d92324', 'get_pptx_title_font_transition__7950498f55dec024fcd49c95968dd860_qw35sft2_a7a84401', 'get_pptx_text_color__f2292754a9e180ffad3521e45a75c905_qw35sft2_d84f0cfa', 'get_slide2_bg_title_bold__3ee304ba526b814e3d4d6c2be946eae0_qw35sft2_71d9d6fc', 'get_pptx_single_text_color__c8d6e08f7da191ce71fc333f5a053d63_qw35sft2_f61679a0', 'get_pptx_all_text_color__2c3419decc4324f20da8cbd8ae0e6c2d_qw35sft2_90378413', 'get_pptx_content_bold__87b69c70091c302ea2d6ddb5f5d9c002_qw35sft2_eb05701d', 'get_impress_slide_pane__0dbd76d4f35e1d6b60aa20d8163a50c2_qw35sft2_d838fcea', 'get_pptx_text_font_size__8cd3b2c6f2a2d80e7aa650d3b0962695_qw35sft2_7904b228', 'get_impress_pptx_props__e526a93f073489d810264dc88333d90d_qw35sft2_38da5688', 'get_pptx_transition__24d25146d014952eed8e287afb1ec612_qw35sft2_f90ee5a0', 'get_pptx_transitions__b146bd229b2aa1513d722fe51dc492dc_qw35sft2_9710b95b', 'get_pptx_table_row0__a30c5776ecd0bc922b282656ddbb8d0e_qw35sft2_1f43e00d', 'get_pptx_slide_title_text__a33045396be2560e0d47268bd7a68b1f_qw35sft2_80ff8e81', 'get_pptx_slide_bg_color__3e8a909d4fdee8dfddabe4251505a445_qw35sft2_a004fe09', 'get_impress_slide_info__b298657a7c2d705873344441c4e6373c_qw35sft2_3775b7a2', 'get_pptx_slide_image_and_title__68b16c8fa261fe6f7f7bf99cd53e207c_qw35sft2_7de26c83', 'get_impress_slide2_title_props__8ff66c4d1cd7e4775354b5fdf25e12eb_qw35sft2_fc2480c5', 'get_pptx_notes__e6eecf5158b3e658cb4aa48b9b009898_qw35sft2_40422256', 'get_pptx_two_slide_tables__422453ea9133a41f5d3b73c5d61c25f4_qw35sft2_68f346c8', 'get_pptx_slide5_text_colors__7491f6c65f733a38f893bdf8cad068b0_qw35sft2_48da877c', 'get_impress_has_audio__5688a332db768af3543ab15e11555c76_qw35sft2_e375e116', 'get_pptx_multi_slide_bg_colors__01d7b3b210af327bea300a303c9bddd6_qw35sft2_eb699e59', 'get_pptx_stretch_and_text__04e2292403d24fd88b4576e3fc170228_qw35sft2_39c2c38f', 'get_lecture_slides_files__ac6b45f9bcb788e83171ced5593fc0f6_qw35sft2_8881d264', 'get_impress_video_docs_vlc__bae283564a19c9e043588bc4818877d3_qw35sft2_2a121f76', 'get_impress_compose4__c60351f10d67b0cf97ff8e001ca27d87_qw35sft2_5f92eccc', 'get_pptx_slide2_bg_and_count__24ae1553823b67948a98514a9ed74ad3_qw35sft2_e8ecc775', 'get_lecture_slides_files__6f1f6989222847d5b9f522129093f325_qw35sft2_5e98a87d', 'get_pptx_notes_and_docx__1f2059f6b413a703d90fb8946ec0e698_qw35sft2_7af43bbc', 'get_pptx_slide2_bg__d5c6210d63b8dbe359320c4bfae3310a_qw35sft2_ed30e2b7', 'get_impress_compose0__9880fcd0ac1db0e1d57c564194f1780a_qw35sft2_622ef45d']

def get_impress_title_color__9e6f539665d909e9d4b70212d289e0ee(env, config: dict):
    """Get title font color from slide 2 of a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/164_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 1)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    color_rgb = None
                                    if run.font.color and run.font.color.rgb:
                                        color_rgb = str(run.font.color.rgb)
                                    return {'title_text': shape.text.strip(), 'color_rgb': color_rgb}
        return {'error': 'Title not found on slide'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__ebdab037a8eab6c57e56f5796ba7c0ab(env, config: dict):
    """Get the title text of a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = ''
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                title_text = shape.text_frame.text.strip()
                break
        return {'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_multi_slide_images__e6e9ce905cbd4fd74addd1fc2cd640de(env, config: dict):
    """Get image properties from multiple slides."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/31_2.pptx')
    slide_indices = config.get('slide_indices', [0, 1])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_data = {}
        for idx in slide_indices:
            if idx >= len(prs.slides):
                slides_data[str(idx)] = {'error': f'Slide index {idx} out of range'}
                continue
            slide = prs.slides[idx]
            images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({'width_emu': shape.width, 'height_emu': shape.height, 'width_cm': round(shape.width / 360000, 2), 'height_cm': round(shape.height / 360000, 2), 'name': shape.name})
            slides_data[str(idx)] = {'image_count': len(images), 'images': images}
        return {'slides': slides_data}
    finally:
        os.unlink(tmp_path)

def get_impress_save_settings__a73132a8bbdfe6ee96f72881e60f8547(env, config: dict):
    """Get auto-save interval and backup copy settings from LibreOffice config."""
    file_path = config.get('path', '/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False, mode='wb') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            content = f.read()
        result = {}
        auto_save_match = re.search('oor:path="/org\\.openoffice\\.Office\\.Common/Save/Document"[^>]*>\\s*<prop\\s+oor:name="AutoSave"[^>]*>\\s*<value>(\\w+)</value>', content)
        result['auto_save_enabled'] = auto_save_match.group(1).lower() == 'true' if auto_save_match else False
        interval_match = re.search('oor:path="/org\\.openoffice\\.Office\\.Common/Save/Document"[^>]*>\\s*<prop\\s+oor:name="AutoSaveTimeIntervall"[^>]*>\\s*<value>(\\d+)</value>', content)
        result['auto_save_minutes'] = int(interval_match.group(1)) if interval_match else None
        backup_match = re.search('oor:path="/org\\.openoffice\\.Office\\.Common/Save/Document"[^>]*>\\s*<prop\\s+oor:name="CreateBackup"[^>]*>\\s*<value>(\\w+)</value>', content)
        result['backup_enabled'] = backup_match.group(1).lower() == 'true' if backup_match else False
        return result
    finally:
        os.unlink(tmp_path)

def get_impress_title_bold__675486c95c0c9441584aa662f56cd677(env, config: dict):
    """Get the title text bold status from a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/4_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        bold_states = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    bold_states.append(bool(run.font.bold))
        return {'title_text': title_shape.text_frame.text.strip(), 'bold_states': bold_states, 'all_bold': all(bold_states) if bold_states else False}
    finally:
        os.unlink(tmp_path)

def get_slide_image_count__07113500ab2edd0ce92ed5a5e9e8b1aa(env, config: dict):
    """Count images on a specific slide of a PPTX file, categorized by size."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 1)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide {slide_index} not found'}
        slide = prs.slides[slide_index]
        images = []
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                images.append({'name': shape.name, 'size': len(shape.image.blob), 'content_type': shape.image.content_type})
        large_images = [img for img in images if img['size'] > 100000]
        return {'total_images': len(images), 'large_images': len(large_images), 'image_details': images}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_rows__d7904b6a5ae76e528f93447180e8419d(env, config: dict):
    """Get the number of rows in a table on a specific slide of a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                row_count = len(table.rows)
                col_count = len(table.columns)
                last_row_values = []
                if row_count > 0:
                    last_row = table.rows[row_count - 1]
                    for i in range(col_count):
                        cell = table.cell(row_count - 1, i)
                        last_row_values.append(cell.text.strip())
                return {'row_count': row_count, 'col_count': col_count, 'last_row_values': last_row_values}
        return {'error': 'No table found on slide'}
    finally:
        os.unlink(tmp_path)

def get_pptx_shape_dimensions__30ab1f17976bddf578367b54740da469(env, config: dict):
    """Get dimensions of specified shapes from a PPTX file."""
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_config = config.get('slides', [])
        results = {}
        for sc in slides_config:
            slide_idx = sc['slide_index']
            shape_name = sc['shape_name']
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.name == shape_name:
                    height_cm = round(shape.height / 360000, 2)
                    width_cm = round(shape.width / 360000, 2)
                    results[f'slide_{slide_idx + 1}_height'] = height_cm
                    results[f'slide_{slide_idx + 1}_width'] = width_cm
                    break
        return results
    finally:
        os.unlink(tmp_path)

def get_impress_slide_table_info__dabff5cb0982227a6c01462d5c667c85(env, config: dict):
    """Get table info (rows, cols) from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = len(list(tbl.rows))
                cols = len(list(tbl.columns))
                tables.append({'rows': rows, 'cols': cols})
        return {'tables': tables, 'table_count': len(tables)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_count__3e661325ce185501287a57bdd6c09b52(env, config: dict):
    """Get slide count and first slide title from a PPTX file."""
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/trimmed.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        first_slide_text = ''
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    first_slide_text = shape.text_frame.text[:200]
                    if first_slide_text.strip():
                        break
        return {'exists': True, 'slide_count': slide_count, 'first_slide_text': first_slide_text.strip()}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__95b2513b7f5b16d49b143853ab08b964(env, config: dict):
    """Get the title text of a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Writing-Outlines.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                title_text = shape.text_frame.text.strip()
                break
        return {'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_impress_title_underline__b1d7e7e1d48378504a47d3ccfb3cf1be(env, config: dict):
    """Get the title text underline status from a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/4_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        underline_states = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    underline_states.append(bool(run.font.underline))
        return {'title_text': title_shape.text_frame.text.strip(), 'underline_states': underline_states, 'all_underlined': all(underline_states) if underline_states else False}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title_and_bg__93a38408d9330c4f554c19c28759d1cc(env, config: dict):
    """Get the title text and background color of a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title_text = text
                    break
        bg_color = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                color = fill.fore_color
                if color and color.rgb:
                    bg_color = str(color.rgb).upper()
            except Exception:
                pass
        return {'title_text': title_text, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_title__497f89854667fa402a3b0a0e21fc061d(env, config: dict):
    """Get title text and alignment from a specific slide in a PPTX file."""
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/22_6.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                text = shape.text_frame.text.strip()
                alignment = None
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.alignment == PP_ALIGN.CENTER:
                        alignment = 'center'
                    elif para.alignment == PP_ALIGN.RIGHT:
                        alignment = 'right'
                    elif para.alignment == PP_ALIGN.LEFT:
                        alignment = 'left'
                    else:
                        alignment = 'left'
                return {'text': text, 'alignment': alignment}
        return {'error': 'Title placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_alignment__d36eb73e90de91fc69ee1df98b9d5058(env, config: dict):
    """Get paragraph alignment of a textbox matching text_pattern on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        text_pattern = config.get('text_pattern', '')
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        align_map = {PP_ALIGN.LEFT: 'LEFT', PP_ALIGN.CENTER: 'CENTER', PP_ALIGN.RIGHT: 'RIGHT', PP_ALIGN.JUSTIFY: 'JUSTIFY'}
        for shape in slide.shapes:
            if shape.has_text_frame and text_pattern.lower() in shape.text_frame.text.lower():
                alignments = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        align_val = para.alignment
                        align_str = align_map.get(align_val, 'NONE')
                        alignments.append(align_str)
                return {'alignments': alignments, 'shape_name': shape.name, 'text': shape.text_frame.text[:100]}
        return {'error': f'No shape with text "{text_pattern}" on slide {slide_index + 1}'}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_exists__ffa5734cd742b977903cbddf387ddf6c(env, config: dict):
    """Check if a table exists on a specific slide of a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        table_count = 0
        for shape in slide.shapes:
            if shape.has_table:
                table_count += 1
        return {'table_exists': table_count > 0, 'table_count': table_count}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title_text__661ea64646d6a38833b5fdf5450470f8(env, config: dict):
    """Get the title text of a specified slide from a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and 'PlaceHolder 1' in shape.name:
                return {'title_text': shape.text.strip()}
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:
                    return {'title_text': shape.text.strip() if shape.has_text_frame else ''}
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return {'title_text': shape.text.strip()}
        return {'title_text': ''}
    finally:
        os.unlink(tmp_path)

def get_impress_middle_text_font_size__e0ab347b041b8b1f40c9ca0e5bc35d39(env, config: dict):
    """Get the font size of the middle textbox (LAUNCH) on slide 1."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/45_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes.append((shape.top, shape))
        text_shapes.sort(key=lambda x: x[0])
        if len(text_shapes) < 2:
            return {'error': 'Less than 2 text shapes found on slide 1'}
        middle_shape = text_shapes[1][1]
        font_sizes = []
        for para in middle_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    font_sizes.append(run.font.size)
        if not font_sizes:
            return {'error': 'No font size found', 'text': middle_shape.text_frame.text}
        size_pt = round(font_sizes[0] / 12700, 1)
        return {'font_size_pt': size_pt, 'text': middle_shape.text_frame.text.strip()}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_bg_color__147ec12b6038f2d22f20700395a41639(env, config: dict):
    """Get the background fill color of a specified slide."""
    from pptx import Presentation
    from pptx.enum.dml import MSO_THEME_COLOR
    file_path = config.get('path', '/home/user/Desktop/71_6.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                color_rgb = str(fill.fore_color.rgb)
                return {'color_rgb': color_rgb}
            except Exception:
                pass
        return {'color_rgb': 'NONE'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_audio__1b19ca869f4a35a02b8d110c7e17b7e9(env, config: dict):
    """Extract audio from a specific slide and compute its MD5 hash.
    Also compute MD5 hash of the reference audio file on the VM.
    Returns both hashes for comparison.
    """
    from pptx import Presentation
    pptx_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    reference_path = config.get('reference_path', '')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'Presentation file not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range (total: {len(prs.slides)})'}
        slide = prs.slides[slide_index]
        audio_data = None
        audio_content_type = None
        for rel in slide.part.rels.values():
            ct = getattr(rel.target_part, 'content_type', '') if hasattr(rel, 'target_part') else ''
            if 'audio' in ct.lower() or ct.lower().startswith('audio/'):
                audio_data = rel.target_part.blob
                audio_content_type = ct
                break
        if audio_data is None:
            for rel in slide.part.rels.values():
                target_ref = rel.target_ref if hasattr(rel, 'target_ref') else ''
                if isinstance(target_ref, str) and ('audio' in target_ref.lower() or '.mp3' in target_ref.lower() or '.wav' in target_ref.lower()):
                    try:
                        audio_data = rel.target_part.blob
                        break
                    except Exception:
                        continue
        if audio_data is None:
            return {'error': 'No audio found on the specified slide', 'has_audio': False}
        slide_audio_hash = hashlib.md5(audio_data).hexdigest()
        ref_bytes = env.controller.get_file(reference_path)
        if not ref_bytes:
            return {'has_audio': True, 'slide_audio_hash': slide_audio_hash, 'error': 'Reference audio file not found'}
        ref_audio_hash = hashlib.md5(ref_bytes).hexdigest()
        return {'has_audio': True, 'slide_audio_hash': slide_audio_hash, 'ref_audio_hash': ref_audio_hash, 'audio_match': slide_audio_hash == ref_audio_hash, 'audio_size': len(audio_data)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__36872752e18a684ac0b72ebb0b9e2c33(env, config: dict):
    """Get the title text from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == 'PlaceHolder 1' and shape.has_text_frame:
                title_text = shape.text_frame.text.strip()
                return {'title_text': title_text}
        return {'title_text': ''}
    finally:
        os.unlink(tmp_path)

def get_impress_table_bg__9da3068d76a130c43b0820c0580d6baa(env, config: dict):
    """Get table cell value and slide background color from a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/181_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        cell_value = None
        row_idx = config.get('row', 1)
        col_idx = config.get('col', 0)
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                if row_idx < len(table.rows) and col_idx < len(table.columns):
                    cell_value = table.cell(row_idx, col_idx).text.strip()
                break
        bg_color = None
        bg = slide.background
        fill = bg.fill
        try:
            if fill.type is not None:
                fg = fill.fore_color
                if fg and fg.rgb:
                    bg_color = str(fg.rgb)
        except Exception:
            bg_color = None
        return {'cell_value': cell_value, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_impress_content_bold__ca58c5a207cc4d7f5ab415f13aba66b0(env, config: dict):
    """Get content/subtitle text and bold property from a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/9_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        content_text = None
        is_bold = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type is not None and int(ph_type) == 7:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                content_text = para.text.strip()
                                for run in para.runs:
                                    is_bold = run.font.bold
                                break
        return {'content_text': content_text, 'is_bold': is_bold}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_row__d2dd546208d48aeec32a49334e180bbc(env, config: dict):
    """Get values from a specific row of a table on a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    row_index = config.get('row_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                if row_index >= len(table.rows):
                    return {'error': f'Row index {row_index} out of range'}
                row = table.rows[row_index]
                values = [cell.text.strip() for cell in row.cells]
                return {'values': values}
        return {'error': 'No table found on slide'}
    finally:
        os.unlink(tmp_path)

def get_impress_image_position__b7ad7d38cb36c439bba7c5f1f082cb99(env, config: dict):
    """Get the position and size of the first picture on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                center_x = shape.left + shape.width / 2
                center_y = shape.top + shape.height / 2
                return {'center_x': center_x, 'center_y': center_y, 'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'slide_width': slide_width, 'slide_height': slide_height}
        return {'error': 'No picture found on slide'}
    finally:
        os.unlink(tmp_path)

def get_impress_bg_fontsize__c22d8b15c916d1ef420eb42daeb27bf0(env, config: dict):
    """Get slide background color and title font size from a pptx file."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        bg_color = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_color = str(fill.fore_color.rgb)
            except Exception:
                bg_color = None
        font_size_pt = None
        for shape in slide.shapes:
            if hasattr(shape, 'text') and 'Should You Consider Studying Abroad' in shape.text:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                font_size_pt = run.font.size / 12700
                                break
                        if font_size_pt is not None:
                            break
                break
        return {'bg_color': bg_color, 'font_size_pt': font_size_pt}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide1_title_format__7d9f2200326202f628d33a5691960f9c(env, config: dict):
    """Get font formatting of the title text on slide 1."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/39_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame and 'Business Infographic Template' in shape.text:
                results = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        italic = run.font.italic
                        color_rgb = None
                        if run.font.color and run.font.color.rgb:
                            color_rgb = str(run.font.color.rgb)
                        results.append({'text': run.text, 'italic': italic, 'color_rgb': color_rgb})
                return {'runs': results}
        return {'error': 'Title text not found on slide 1'}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_bold__c4941aca377cb2eab0e421e15f393854(env, config: dict):
    """Get bold status of text runs in a specific shape of a pptx slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/69_4.pptx')
    slide_index = config.get('slide_index', 0)
    shape_index = config.get('shape_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        runs_bold = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                runs_bold.append({'text': run.text, 'bold': bool(run.font.bold)})
        all_bold = all((r['bold'] for r in runs_bold)) if runs_bold else False
        return {'all_bold': all_bold, 'runs': runs_bold}
    finally:
        os.unlink(tmp_path)

def get_slide_orientation__662bf4c59a75818f259c7e0aef709493(env, config: dict):
    """Get slide width and height from a pptx file to determine orientation."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        width = prs.slide_width
        height = prs.slide_height
        return {'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_font_size__a1c0faa26bf29b13cb4b3dfa5122d14a(env, config: dict):
    """Get font size of title text in a pptx slide."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/69_4.pptx')
    slide_index = config.get('slide_index', 0)
    shape_index = config.get('shape_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        font_sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    font_sizes.append(run.font.size / Pt(1))
        return {'font_sizes': font_sizes}
    finally:
        os.unlink(tmp_path)

def get_impress_title_style_bg__d816da51e2714a5b86252b84b6e77125(env, config: dict):
    """Get title bold/italic status and slide background color from pptx."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    file_path = config.get('path', '/home/user/Desktop/16_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        title_shape = slide.shapes[0]
        title_bold = None
        title_italic = None
        if title_shape.has_text_frame:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    title_bold = run.font.bold
                    title_italic = run.font.italic
                    break
                break
        bg_color = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_color = str(fill.fore_color.rgb)
            except Exception:
                bg_color = None
        return {'title_bold': title_bold, 'title_italic': title_italic, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_slide1_image_position__ebf21364e626ee60bab03f3d9d12deff(env, config: dict):
    """Get position of images on slide 1 to check if any image has been moved to target position."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        images_info = []
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                try:
                    w_cm = shape.width / 360000
                    h_cm = shape.height / 360000
                    x_cm = shape.left / 360000
                    y_cm = shape.top / 360000
                    area = w_cm * h_cm
                    images_info.append({'name': shape.name, 'width_cm': round(w_cm, 2), 'height_cm': round(h_cm, 2), 'left_cm': round(x_cm, 2), 'top_cm': round(y_cm, 2), 'area_cm2': round(area, 2)})
                except Exception:
                    pass
        images_info.sort(key=lambda x: x.get('area_cm2', 0), reverse=True)
        return {'images': images_info}
    finally:
        os.unlink(tmp_path)

def get_impress_subtitle_font_color__d0c681af3bd2af47798202cb5ebe1988(env, config: dict):
    """Get the font color of the subtitle (content placeholder) text on a specified slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/71_6.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 2':
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            return {'color_rgb': str(run.font.color.rgb)}
                return {'error': 'No runs with color found in subtitle'}
        return {'error': 'Subtitle placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_slide_text_font_color__3892513c1cab5055660dfeaa188121fb(env, config: dict):
    """Get font color of text runs on a specific slide and shape in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/1_2.pptx')
    slide_index = config.get('slide_index', 2)
    shape_index = config.get('shape_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            return {'error': f'Shape index {shape_index} out of range'}
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        colors = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    colors.append(str(run.font.color.rgb))
                else:
                    colors.append(None)
        return {'colors': colors, 'text': shape.text_frame.text}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__7b96cb41298a29a8e5448750424ab134(env, config: dict):
    """Get the title text of a specific slide from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_text = shape.text_frame.text.strip()
                break
        return {'title': title_text, 'slide_count': len(prs.slides)}
    finally:
        os.unlink(tmp_path)

def get_impress_strikethrough__a345b8271c8198e8c95ec913d78c781f(env, config: dict):
    """Get strikethrough status of all paragraphs on a specific slide."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        paragraphs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    all_strike = True
                    has_runs = False
                    for run in para.runs:
                        if run.text.strip():
                            has_runs = True
                            if not run.font.strikethrough:
                                all_strike = False
                                break
                    paragraphs.append({'text': text, 'strikethrough': all_strike and has_runs})
        return {'paragraphs': paragraphs}
    finally:
        os.unlink(tmp_path)

def get_pptx_pic_width_and_fonts__1a124508c5c52deecd8c9cfabb1cdd16(env, config: dict):
    """Get picture width on a target slide and font sizes of textboxes on another slide."""
    import openpyxl
    from pptx import Presentation
    from pptx.util import Emu, Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/42_2.pptx')
    pic_slide_index = config.get('pic_slide_index', 2)
    font_slide_index = config.get('font_slide_index', 3)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        pic_width_cm = None
        if pic_slide_index < len(slides):
            slide = slides[pic_slide_index]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic_width_cm = round(shape.width / 360000, 2)
                    break
        font_sizes = []
        if font_slide_index < len(slides):
            slide = slides[font_slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                font_sizes.append(round(run.font.size / 12700, 1))
        return {'pic_width_cm': pic_width_cm, 'font_sizes': font_sizes}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_info__b84f4f9ed96e0e954dfada4f0a280734(env, config: dict):
    """Get slide count and first text of specified slide positions from a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide_count = len(slides)
        check_indices = config.get('check_indices', [0])
        slide_texts = {}
        for idx in check_indices:
            actual_idx = idx
            if idx < 0:
                actual_idx = slide_count + idx
            if 0 <= actual_idx < slide_count:
                texts = []
                for shape in slides[actual_idx].shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                                break
                        if texts:
                            break
                slide_texts[str(idx)] = texts[0] if texts else ''
            else:
                slide_texts[str(idx)] = None
        return {'slide_count': slide_count, 'slide_texts': slide_texts}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_font_color__3536f82770e6c6d12bb7de4bc43120d5(env, config: dict):
    """Get the font color of the title text on a specified slide."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/21_0.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        colors = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    colors.append(str(run.font.color.rgb))
                else:
                    colors.append(None)
        return {'title_text': title_shape.text_frame.text, 'font_colors': colors, 'first_color': colors[0] if colors else None}
    finally:
        os.unlink(tmp_path)

def get_pptx_image_props__32c4d5ced54cde2e00ee07b6ba3973e4(env, config: dict):
    """Get image properties (count, width, height) from a specific slide."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/31_2.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        images = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({'width_emu': shape.width, 'height_emu': shape.height, 'width_cm': round(shape.width / 360000, 2), 'height_cm': round(shape.height / 360000, 2), 'name': shape.name})
        return {'image_count': len(images), 'images': images}
    finally:
        os.unlink(tmp_path)

def get_impress_bg_and_subtitle__024a68353adff6dd2d6a01cc842a46ed(env, config: dict):
    """Get background colors of slides 1 & 6, and subtitle text of slide 2."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/13_0.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide1_bg = None
        if len(slides) >= 1:
            fill = slides[0].background.fill
            if fill.type is not None:
                try:
                    slide1_bg = str(fill.fore_color.rgb)
                except Exception:
                    pass
        slide6_bg = None
        if len(slides) >= 6:
            fill = slides[5].background.fill
            if fill.type is not None:
                try:
                    slide6_bg = str(fill.fore_color.rgb)
                except Exception:
                    pass
        slide2_subtitle = None
        if len(slides) >= 2:
            slide2 = slides[1]
            text_shapes = []
            for shape in slide2.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and para.runs:
                            fsize = para.runs[0].font.size or 0
                            text_shapes.append((fsize, text))
            text_shapes.sort(key=lambda x: x[0])
            if text_shapes:
                slide2_subtitle = text_shapes[0][1]
        return {'slide1_bg_color': slide1_bg, 'slide6_bg_color': slide6_bg, 'slide2_subtitle': slide2_subtitle}
    finally:
        os.unlink(tmp_path)

def get_slide_title_font_color__66374a642830d11922d74321b71d53a4(env, config: dict):
    """Get the font color of the title text on a specific slide."""
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 0:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                rgb_str = str(rgb)
                                r = int(rgb_str[0:2], 16)
                                g = int(rgb_str[2:4], 16)
                                b = int(rgb_str[4:6], 16)
                                return {'color': [r, g, b], 'color_hex': rgb_str, 'text': run.text}
                    return {'error': 'Title has no explicit font color'}
            except (ValueError, AttributeError):
                continue
        return {'error': 'Title placeholder not found on this slide'}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_format__322bfa5ed10c48cee5b1ded75b68bcd6(env, config: dict):
    """Get the title font formatting of a specified slide from a pptx file."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and 'PlaceHolder 1' in shape.name:
                result = {'title_text': shape.text.strip(), 'bold': False, 'font_color_rgb': None, 'font_size_pt': None}
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        result['bold'] = bool(run.font.bold)
                        if run.font.color and run.font.color.rgb:
                            result['font_color_rgb'] = str(run.font.color.rgb)
                        if run.font.size:
                            result['font_size_pt'] = run.font.size / 12700
                        break
                    break
                return result
        return {'error': 'Title placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_pptx_shape_dimensions__37468b49cfdc53eaca3dfbbee86985a3(env, config: dict):
    """Get dimensions of specified shapes from a PPTX file."""
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_config = config.get('slides', [])
        results = {}
        for sc in slides_config:
            slide_idx = sc['slide_index']
            shape_name = sc['shape_name']
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.name == shape_name:
                    height_cm = round(shape.height / 360000, 2)
                    width_cm = round(shape.width / 360000, 2)
                    results[f'slide_{slide_idx + 1}_height'] = height_cm
                    results[f'slide_{slide_idx + 1}_width'] = width_cm
                    break
        return results
    finally:
        os.unlink(tmp_path)

def get_pptx_title_text__927f9acbea3f73fbaa9e281542feb428(env, config: dict):
    """Get the title text on a specified slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/21_0.pptx')
    slide_index = config.get('slide_index', 3)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        return {'title_text': title_shape.text_frame.text.strip()}
    finally:
        os.unlink(tmp_path)

def get_impress_content_align_title_color__63833ad040c5e27e84bd66675fcf45fe(env, config: dict):
    """Get content text alignment and title font color from pptx."""
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/16_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        title_shape = slide.shapes[0]
        title_color = None
        if title_shape.has_text_frame:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and run.font.color.type:
                        title_color = str(run.font.color.rgb)
                    break
                break
        content_shape = slide.shapes[1]
        content_alignment = None
        if content_shape.has_text_frame:
            for para in content_shape.text_frame.paragraphs:
                if para.alignment is not None:
                    content_alignment = str(para.alignment)
                break
        return {'title_color': title_color, 'content_alignment': content_alignment}
    finally:
        os.unlink(tmp_path)

def get_slide_count__55d37620fcf7998461b41e2b1579f16b(env, config: dict):
    """Get the number of slides in a pptx file."""
    import tempfile
    import os
    from pptx import Presentation
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        count = len(prs.slides)
        return {'slide_count': count}
    finally:
        os.unlink(tmp_path)

def get_impress_title_bold__6e02436ff918c1445d486be48242a3ca(env, config: dict):
    """Get bold property of a specific slide title from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/24_8.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        return {'bold': bool(run.font.bold), 'text': run.text, 'slide_index': slide_index}
        return {'error': 'No title found'}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_info__1abdd5c6fe686f66f2606ee0a55bf0d1(env, config: dict):
    """Get slide count and first text of specified slide positions from a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide_count = len(slides)
        check_indices = config.get('check_indices', [0])
        slide_texts = {}
        for idx in check_indices:
            if 0 <= idx < slide_count:
                texts = []
                for shape in slides[idx].shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                                break
                        if texts:
                            break
                slide_texts[str(idx)] = texts[0] if texts else ''
            else:
                slide_texts[str(idx)] = None
        return {'slide_count': slide_count, 'slide_texts': slide_texts}
    finally:
        os.unlink(tmp_path)

def get_pptx_textbox_font_sizes__a0ffb78fbe8b896461c46fcb419a2efb(env, config: dict):
    """Get font sizes of textboxes on a specific slide by matching text content."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    target_texts = config.get('target_texts', [])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        results = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text in target_texts:
                        sizes = set()
                        for run in para.runs:
                            if run.font.size:
                                sizes.add(round(run.font.size / 12700, 1))
                        if sizes:
                            results[text] = list(sizes)[0]
        return {'font_sizes': results}
    finally:
        os.unlink(tmp_path)

def get_impress_all_title_fonts__9c7bd04930deae5cddd57fa6475996f4(env, config: dict):
    """Get font names of all slide titles from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/24_8.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        fonts = {}
        for (i, slide) in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            fonts[f'slide_{i + 1}'] = run.font.name
                            break
                    break
        return {'fonts': fonts, 'slide_count': len(prs.slides)}
    finally:
        os.unlink(tmp_path)

def get_slide_count_and_layout__c9d9670d65b42979a3ba7969e086139a(env, config: dict):
    """Get slide count and check if all slides are blank (no text shapes)."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        slides_with_text = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slides_with_text += 1
                    break
        return {'slide_count': slide_count, 'slides_with_text': slides_with_text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_text_and_bg__6d65cdf6ba3aad0bcf7dd7e1b9c52c8f(env, config: dict):
    """Get slide 6 main text and slide 3 background color from the presentation."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/13_0.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide6_text = None
        max_font_size = 0
        if len(slides) >= 6:
            slide6 = slides[5]
            for shape in slide6.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and para.runs:
                            run = para.runs[0]
                            fsize = run.font.size or 0
                            if fsize > max_font_size:
                                max_font_size = fsize
                                slide6_text = text
        slide3_bg = None
        if len(slides) >= 3:
            slide3 = slides[2]
            fill = slide3.background.fill
            if fill.type is not None:
                try:
                    slide3_bg = str(fill.fore_color.rgb)
                except Exception:
                    slide3_bg = None
        return {'slide6_text': slide6_text, 'slide3_bg_color': slide3_bg}
    finally:
        os.unlink(tmp_path)

def get_impress_title_fontsize__f67d7e69d6a85882636e0a41dabd5c89(env, config: dict):
    """Get font size of a specific slide title from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/24_8.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        size_emu = run.font.size
                        size_pt = size_emu / 12700 if size_emu else None
                        return {'font_size_pt': size_pt, 'font_size_emu': size_emu, 'text': run.text, 'slide_index': slide_index}
        return {'error': 'No title found'}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_underline__1563ed9d69890bfb1ca48bd048aeee4a(env, config: dict):
    """Get underline status of text runs in a specific shape of a pptx slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/69_4.pptx')
    slide_index = config.get('slide_index', 0)
    shape_index = config.get('shape_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        runs_underline = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                runs_underline.append({'text': run.text, 'underline': bool(run.font.underline)})
        all_underline = all((r['underline'] for r in runs_underline)) if runs_underline else False
        return {'all_underline': all_underline, 'runs': runs_underline}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_bold__1f6e651edc27326d3dac73dc9e29333b(env, config: dict):
    """Check if all text in the table has bold formatting."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/154_3.pptx')
    table_name = config.get('table_name', 'Table 3')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == table_name and shape.has_table:
                table = shape.table
                total_runs = 0
                bold_runs = 0
                for row in table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    total_runs += 1
                                    if run.font.bold:
                                        bold_runs += 1
                return {'total_runs': total_runs, 'bold_runs': bold_runs, 'all_bold': total_runs > 0 and bold_runs == total_runs}
        return {'error': f"Table '{table_name}' not found"}
    finally:
        os.unlink(tmp_path)

def get_pptx_textbox_font_sizes__8e211b70307c1527c1dc374c2a15baa9(env, config: dict):
    """Get font sizes of textboxes on a specific slide by matching text content."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    target_texts = config.get('target_texts', [])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        results = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text in target_texts:
                        sizes = set()
                        for run in para.runs:
                            if run.font.size:
                                sizes.add(round(run.font.size / 12700, 1))
                        if sizes:
                            results[text] = list(sizes)[0]
        return {'font_sizes': results}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_alignment__b3837f805dc03cde831accc6a7063290(env, config: dict):
    """Get paragraph alignment of a textbox matching text_pattern on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        text_pattern = config.get('text_pattern', '')
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        align_map = {PP_ALIGN.LEFT: 'LEFT', PP_ALIGN.CENTER: 'CENTER', PP_ALIGN.RIGHT: 'RIGHT', PP_ALIGN.JUSTIFY: 'JUSTIFY'}
        for shape in slide.shapes:
            if shape.has_text_frame and text_pattern.lower() in shape.text_frame.text.lower():
                alignments = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        align_val = para.alignment
                        align_str = align_map.get(align_val, 'NONE')
                        alignments.append(align_str)
                return {'alignments': alignments, 'shape_name': shape.name, 'text': shape.text_frame.text[:100]}
        return {'error': f'No shape with text "{text_pattern}" on slide {slide_index + 1}'}
    finally:
        os.unlink(tmp_path)

def get_pptx_orientation__2f288a1714ed5a4b53fed9b568a61679(env, config: dict):
    """Get slide orientation from a PPTX file."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        width = prs.slide_width
        height = prs.slide_height
        orientation = 'portrait' if height > width else 'landscape'
        return {'orientation': orientation, 'width': int(width), 'height': int(height)}
    finally:
        os.unlink(tmp_path)

def get_pptx_bold_check__c8d3ea84bab1df08b4dcc2adb9177ac5(env, config: dict):
    """Check if text shapes on specified slides have bold formatting."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/43_1.pptx')
    slide_indices = config.get('slide_indices', [0, 1])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        results = {}
        for idx in slide_indices:
            if idx >= len(prs.slides):
                results[f'slide_{idx}'] = {'error': 'Slide not found'}
                continue
            slide = prs.slides[idx]
            bold_runs = 0
            total_runs = 0
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text.strip():
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                total_runs += 1
                                if run.font.bold:
                                    bold_runs += 1
            results[f'slide_{idx}'] = {'bold_runs': bold_runs, 'total_runs': total_runs}
        return results
    finally:
        os.unlink(tmp_path)

def get_pptx_textbox_font_sizes__d842f0d6650bc2853fbd82b34d9bb662(env, config: dict):
    """Get font sizes of textboxes on a specific slide by matching text content."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    target_texts = config.get('target_texts', [])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        results = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text in target_texts:
                        sizes = set()
                        for run in para.runs:
                            if run.font.size:
                                sizes.add(round(run.font.size / 12700, 1))
                        if sizes:
                            results[text] = list(sizes)[0]
        return {'font_sizes': results}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_info__e9203909bdddcf8283496218733a0b9b(env, config: dict):
    """Get slide count and first two slide titles from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Forests.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        titles = []
        for idx in range(min(slide_count, 2)):
            title = ''
            for shape in prs.slides[idx].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            title = text
                            break
                if title:
                    break
            titles.append(title)
        return {'slide_count': slide_count, 'slide1_title': titles[0] if len(titles) > 0 else '', 'slide2_title': titles[1] if len(titles) > 1 else ''}
    finally:
        os.unlink(tmp_path)

def get_slide_notes__44fd5b3a7c457950926c4afa82f02524(env, config: dict):
    """Get speaker notes text from a specific slide of a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        if slide_index >= len(slides):
            return {'error': f'Slide index {slide_index} out of range (total: {len(slides)})'}
        slide = slides[slide_index]
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
        return {'notes_text': notes_text.strip(), 'slide_index': slide_index}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__07d639643e02063bb060c675d5eb936b(env, config: dict):
    """Get the title text from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and title_text is None:
                    title_text = text
        return {'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide2_welcome_format__203dc282bb468063e328a7afef0f556d(env, config: dict):
    """Get font formatting of the welcome message on slide 2."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/39_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[1]
        for shape in slide.shapes:
            if shape.has_text_frame and 'Hello & welcome' in shape.text:
                results = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        size_pt = None
                        if run.font.size:
                            size_pt = run.font.size.pt
                        results.append({'text': run.text, 'underline': run.font.underline, 'size_pt': size_pt})
                return {'runs': results}
        return {'error': 'Welcome message not found on slide 2'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_text__64fdb0e52d9e147a698cbe6791191c81(env, config: dict):
    """Get all text content from a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    target_shape_name = config.get('shape_name', None)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if target_shape_name and shape.name != target_shape_name:
                    continue
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        return {'texts': texts}
    finally:
        os.unlink(tmp_path)

def get_impress_bg_font_color__d8eda1446480065d3fccc5f3a4fe966e(env, config: dict):
    """Get slide background color and title font color from a pptx file."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    file_path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        bg_color = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_color = str(fill.fore_color.rgb)
            except Exception:
                bg_color = None
        font_color = None
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                font_color = str(run.font.color.rgb)
                                break
                        if font_color:
                            break
            if font_color:
                break
        return {'bg_color': bg_color, 'font_color': font_color}
    finally:
        os.unlink(tmp_path)

def get_impress_top_text_color__42394c3fddd5c54d80cb5da5e81a07d8(env, config: dict):
    """Get the font color of the topmost textbox on slide 1 of a pptx file."""
    from pptx import Presentation
    from pptx.util import Emu
    file_path = config.get('path', '/home/user/Desktop/45_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes.append((shape.top, shape))
        text_shapes.sort(key=lambda x: x[0])
        if not text_shapes:
            return {'error': 'No text shapes found on slide 1'}
        top_shape = text_shapes[0][1]
        colors = []
        for para in top_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    colors.append(str(run.font.color.rgb))
        if not colors:
            return {'error': 'No color found', 'text': top_shape.text_frame.text}
        return {'color': colors[0], 'text': top_shape.text_frame.text.strip(), 'all_colors': colors}
    finally:
        os.unlink(tmp_path)

def get_impress_table_with_content__608b93868302b656e50a31095c9d18e9(env, config: dict):
    """Get table dimensions and first cell content from a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    slide_index = config.get('slide_index', 2)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = len(list(tbl.rows))
                cols = len(list(tbl.columns))
                first_cell_text = tbl.cell(0, 0).text.strip()
                return {'rows': rows, 'cols': cols, 'first_cell_text': first_cell_text, 'has_table': True}
        return {'has_table': False}
    finally:
        os.unlink(tmp_path)

def get_slide_orientation__060b6ae42d4d041f46057b0dff06cb53(env, config: dict):
    """Get slide orientation from a PPTX file."""
    import tempfile
    import os
    path = config.get('path', '/home/user/Downloads/Secrets-of-Monetizing-Video.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        width = prs.slide_width
        height = prs.slide_height
        orientation = 'portrait' if height > width else 'landscape'
        return {'width': width, 'height': height, 'orientation': orientation}
    finally:
        os.unlink(tmp_path)

def get_impress_notes_bg__7bae80483f423bfa3b3bbb7f9e9eafdb(env, config: dict):
    """Get slide notes text and background color from a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/181_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        bg_color = None
        bg = slide.background
        fill = bg.fill
        try:
            if fill.type is not None:
                fg = fill.fore_color
                if fg and fg.rgb:
                    bg_color = str(fg.rgb)
        except Exception:
            bg_color = None
        return {'notes_text': notes_text, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_slide_count__1573725569136f71d4ad79ac5ffbf35a(env, config: dict):
    """Get the number of slides in a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        return {'slide_count': len(prs.slides)}
    finally:
        os.unlink(tmp_path)

def get_slide_bg_color__1dbf2adcf932ab2d0e517ac3d78f9d08(env, config: dict):
    """Get background color of a specific slide in a pptx file."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.dml import MSO_THEME_COLOR
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        bg = slide.background
        fill = bg.fill
        try:
            from pptx.enum.dml import MSO_FILL
            if fill.type == MSO_FILL.SOLID:
                rgb = fill.fore_color.rgb
                hex_str = str(rgb)
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
                return {'r': r, 'g': g, 'b': b, 'has_color': True}
        except Exception:
            pass
        try:
            from lxml import etree
            bg_elem = slide.background._element
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            solid_fill = bg_elem.findall('.//a:solidFill', ns)
            if solid_fill:
                srgb = solid_fill[0].find('a:srgbClr', ns)
                if srgb is not None:
                    hex_val = srgb.get('val')
                    r = int(hex_val[0:2], 16)
                    g = int(hex_val[2:4], 16)
                    b = int(hex_val[4:6], 16)
                    return {'r': r, 'g': g, 'b': b, 'has_color': True}
        except Exception:
            pass
        return {'r': 255, 'g': 255, 'b': 255, 'has_color': False}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_info__88930a069947104f21bb2b61fdebcc2c(env, config: dict):
    """Get slide count and first text of specified slide positions from a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide_count = len(slides)
        check_indices = config.get('check_indices', [0])
        slide_texts = {}
        for idx in check_indices:
            actual_idx = idx
            if idx < 0:
                actual_idx = slide_count + idx
            if 0 <= actual_idx < slide_count:
                texts = []
                for shape in slides[actual_idx].shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                                break
                        if texts:
                            break
                slide_texts[str(idx)] = texts[0] if texts else ''
            else:
                slide_texts[str(idx)] = None
        return {'slide_count': slide_count, 'slide_texts': slide_texts}
    finally:
        os.unlink(tmp_path)

def get_impress_title_color__6b28d17bfe07bb684917211f85bf497c(env, config: dict):
    """Get the title text font color from a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/4_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        colors = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    color_rgb = None
                    if run.font.color and run.font.color.rgb:
                        color_rgb = str(run.font.color.rgb)
                    colors.append(color_rgb)
        return {'title_text': title_shape.text_frame.text.strip(), 'colors': colors, 'all_same_color': len(set(colors)) <= 1, 'primary_color': colors[0] if colors else None}
    finally:
        os.unlink(tmp_path)

def get_slide_image_info__e8973b7aba53f626fa037406b1da4c8e(env, config: dict):
    """Get info about images and text placeholders on the first slide."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(tmp_path)
        if len(prs.slides) == 0:
            return {'error': 'No slides found'}
        slide = prs.slides[0]
        image_count = 0
        text_placeholder_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if shape.has_text_frame and shape.text_frame.text.strip():
                text_placeholder_count += 1
            elif hasattr(shape, 'is_placeholder') and shape.is_placeholder and shape.has_text_frame:
                text_placeholder_count += 1
        return {'image_count': image_count, 'text_placeholder_count': text_placeholder_count}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_title_bg__c2ac2f94e6c641c5a5c2e733ff74ade1(env, config: dict):
    """Get slide title text and background color from a pptx file."""
    from pptx import Presentation
    from pptx.enum.dml import MSO_THEME_COLOR
    file_path = config.get('path', '/home/user/Desktop/181_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        title_text = ''
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                title_text = shape.text_frame.text.strip()
                break
        bg_color = None
        bg = slide.background
        fill = bg.fill
        try:
            if fill.type is not None:
                fg = fill.fore_color
                if fg and fg.rgb:
                    bg_color = str(fg.rgb)
        except Exception:
            bg_color = None
        return {'title_text': title_text, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_paragraphs__10c4d6d7b532dd358d2936a84f08db4a(env, config: dict):
    """Get all paragraph texts from a specific placeholder on a slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Writing-Outlines.pptx')
    slide_index = config.get('slide_index', 0)
    placeholder_name = config.get('placeholder_name', 'PlaceHolder 2')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        paragraphs = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == placeholder_name:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs.append(text)
                break
        return {'paragraphs': paragraphs, 'count': len(paragraphs)}
    finally:
        os.unlink(tmp_path)

def get_snapshot_and_slide_bg__da1a5547fec970e705ee94d1f2bd1e91(env, config: dict):
    """Check both VLC snapshot existence and slide background image."""
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    result = {'snapshot_exists': False, 'has_bg_image': False}
    try:
        snap_result = env.controller.run_bash_script('ls /home/user/Pictures/vlcsnap-*.png 2>/dev/null | head -1', timeout=30)
        output = snap_result.get('output', '').strip() if isinstance(snap_result, dict) else ''
        if output and 'vlcsnap-' in output:
            result['snapshot_exists'] = True
    except Exception as e:
        logger.error('Error checking snapshot: %s', str(e))
    try:
        file_bytes = env.controller.get_file(ppt_file_path)
        if not file_bytes:
            return result
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as myzip:
                slide_xml = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
                if slide_xml not in myzip.namelist():
                    return result
                with myzip.open(slide_xml) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                    image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                    attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                    for child in root.iter(bg_tag):
                        for element in child.iter(image_tag):
                            image_id = element.attrib.get(attr_tag, '')
                            if image_id:
                                result['has_bg_image'] = True
                                break
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error('Error checking slide bg: %s', str(e))
    return result

def get_pptx_slide_dimensions__54270a3c8d106e03ce2d0c72089a3417(env, config: dict):
    """Get slide width and height in cm from a pptx file."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        width_cm = round(prs.slide_width.cm, 2)
        height_cm = round(prs.slide_height.cm, 2)
        return {'width_cm': width_cm, 'height_cm': height_cm}
    finally:
        os.unlink(tmp_path)

def get_impress_bottom_text_bold__06366a1cfa1039e4258f94612f6055b9(env, config: dict):
    """Get the bold property of the bottom textbox on slide 1."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/45_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_shapes.append((shape.top, shape))
        text_shapes.sort(key=lambda x: x[0])
        if len(text_shapes) < 3:
            return {'error': 'Less than 3 text shapes found on slide 1'}
        bottom_shape = text_shapes[2][1]
        bold_states = []
        for para in bottom_shape.text_frame.paragraphs:
            for run in para.runs:
                bold_states.append(run.font.bold)
        return {'bold_states': bold_states, 'all_bold': all((b is True for b in bold_states)), 'text': bottom_shape.text_frame.text.strip()}
    finally:
        os.unlink(tmp_path)

def get_pptx_pic_height_and_fonts__4eb9ebb1f165db3d3eb60748b0824590(env, config: dict):
    """Get picture height on a target slide and font sizes of textboxes on another slide."""
    from pptx import Presentation
    from pptx.util import Emu, Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/42_2.pptx')
    pic_slide_index = config.get('pic_slide_index', 3)
    font_slide_index = config.get('font_slide_index', 5)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        pic_height_cm = None
        if pic_slide_index < len(slides):
            slide = slides[pic_slide_index]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic_height_cm = round(shape.height / 360000, 2)
                    break
        font_sizes = []
        if font_slide_index < len(slides):
            slide = slides[font_slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                font_sizes.append(round(run.font.size / 12700, 1))
        return {'pic_height_cm': pic_height_cm, 'font_sizes': font_sizes}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_text__38a0fe7ffb236efbdbbcf3cc7460e807(env, config: dict):
    """Get text from a specific shape on a specific slide of a PPTX file."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 1)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide {slide_index} not found'}
        slide = prs.slides[slide_index]
        shape_name = config.get('shape_name', 'Rectangle 15')
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                return {'text': shape.text.strip()}
        return {'error': f'Shape {shape_name} not found'}
    finally:
        os.unlink(tmp_path)

def get_slide_audio_info__5688a332db768af3543ab15e11555c76(env, config: dict):
    """Get audio information from a specific slide in a pptx file.

    Config:
        path: path to the pptx file on VM
        slide_index: 0-based slide index to check for audio

    Returns:
        dict with 'has_audio' (bool), 'audio_count' (int), 'slide_count' (int)
    """
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_audio': False}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide_count = len(slides)
        if slide_index < 0 or slide_index >= slide_count:
            return {'error': f'Slide index {slide_index} out of range (0-{slide_count - 1})', 'has_audio': False, 'slide_count': slide_count}
        slide = slides[slide_index]
        audio_count = 0
        audio_names = []
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                audio_count += 1
                audio_names.append(shape.name if hasattr(shape, 'name') else 'unknown')
        if audio_count == 0:
            for rel in slide.part.rels.values():
                if rel.reltype and 'audio' in rel.reltype.lower():
                    audio_count += 1
                elif hasattr(rel, 'target_ref') and rel.target_ref and any((ext in rel.target_ref.lower() for ext in ['.mp3', '.wav', '.ogg', '.m4a'])):
                    audio_count += 1
        return {'has_audio': audio_count > 0, 'audio_count': audio_count, 'slide_count': slide_count, 'slide_index_checked': slide_index}
    except Exception as e:
        return {'error': str(e), 'has_audio': False}
    finally:
        os.unlink(tmp_path)

def get_impress_title_font_props__8c5c9bdd88c0ba0d1f43df42a252cfae(env, config: dict):
    """Get title font properties (italic, size) from slide 2 of a pptx file."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/164_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 1)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0:
                        text = shape.text.strip()
                        italic = None
                        size_pt = None
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.text.strip():
                                    italic = run.font.italic
                                    if run.font.size is not None:
                                        size_pt = run.font.size / 12700
                                    break
                            if italic is not None:
                                break
                        return {'title_text': text, 'italic': italic, 'size_pt': size_pt}
        return {'error': 'Title not found on slide'}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_bold__7728836b00929340fdef9a2a6602ef09(env, config: dict):
    """Get the bold property of the title text on a specified slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/21_0.pptx')
    slide_index = config.get('slide_index', 2)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_shape = None
        max_font_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size > max_font_size:
                            max_font_size = run.font.size
                            title_shape = shape
        if title_shape is None:
            return {'error': 'No title shape found'}
        bold_values = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                bold_values.append(run.font.bold)
        return {'title_text': title_shape.text_frame.text, 'bold_values': bold_values, 'all_bold': all((b is True for b in bold_values))}
    finally:
        os.unlink(tmp_path)

def get_impress_slide2_state__33bd53bce176d3722b0b7a880cb5489e(env, config: dict):
    """Get slide 2 title underline status and subtitle/body text content."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/164_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 1)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_underline = None
        title_text = None
        body_text = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:
                    title_text = shape.text.strip()
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                title_underline = run.font.underline
                                break
                        if title_underline is not None:
                            break
                elif ph_idx == 1:
                    body_text = shape.text.strip()
        return {'title_text': title_text, 'title_underline': title_underline, 'body_text': body_text}
    finally:
        os.unlink(tmp_path)

def get_slide_text_font_color__175bcb8418fdbfeda65db92690d9b128(env, config: dict):
    """Get font color of text runs on a specific slide and shape in a pptx file."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    file_path = config.get('path', '/home/user/Desktop/1_2.pptx')
    slide_index = config.get('slide_index', 4)
    shape_index = config.get('shape_index', 3)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            return {'error': f'Shape index {shape_index} out of range'}
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        colors = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    colors.append(str(run.font.color.rgb))
                else:
                    colors.append(None)
        return {'colors': colors, 'text': shape.text_frame.text}
    finally:
        os.unlink(tmp_path)

def get_slide_text_bold__db0f57f6a04c8947ec0364b499605e29(env, config: dict):
    """Get bold status of text runs on a specific slide and shape in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/1_2.pptx')
    slide_index = config.get('slide_index', 4)
    shape_index = config.get('shape_index', 4)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        if shape_index >= len(slide.shapes):
            return {'error': f'Shape index {shape_index} out of range'}
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        bold_values = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                bold_values.append(bool(run.font.bold))
        return {'bold_values': bold_values, 'text': shape.text_frame.text}
    finally:
        os.unlink(tmp_path)

def get_impress_title_format__110ae328b957703cc4b8db2a2b04df5c(env, config: dict):
    """Get title text formatting (bold, italic) from a specific slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/22_6.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                text = shape.text_frame.text.strip()
                bold = False
                italic = False
                if shape.text_frame.paragraphs:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold:
                                bold = True
                            if run.font.italic:
                                italic = True
                return {'text': text, 'bold': bold, 'italic': italic}
        return {'error': 'Title placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_impress_strikethrough__bbd5669ff0c1c503844620466c04ffda(env, config: dict):
    """Get strikethrough status of all paragraphs on a specific slide."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        slide = prs.slides[slide_index]
        paragraphs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    all_strike = True
                    has_runs = False
                    for run in para.runs:
                        if run.text.strip():
                            has_runs = True
                            if not run.font.strikethrough:
                                all_strike = False
                                break
                    paragraphs.append({'text': text, 'strikethrough': all_strike and has_runs})
        return {'paragraphs': paragraphs}
    finally:
        os.unlink(tmp_path)

def get_slide_notes__4bc42c00ac052a9214daa0b289a72718(env, config: dict):
    """Get speaker notes text from a specific slide of a pptx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        if slide_index >= len(slides):
            return {'error': f'Slide index {slide_index} out of range (total: {len(slides)})'}
        slide = slides[slide_index]
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
        return {'notes_text': notes_text.strip(), 'slide_index': slide_index}
    finally:
        os.unlink(tmp_path)

def get_pptx_last_slide_info__964158cb23267c07918d09a84fcd28cc(env, config: dict):
    """Get info about the last slide of a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        last_slide = prs.slides[-1]
        last_slide_title = None
        for shape in last_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                last_slide_title = shape.text_frame.text.strip()
                break
        return {'slide_count': slide_count, 'last_slide_title': last_slide_title}
    finally:
        os.unlink(tmp_path)

def get_slide_orientation__326a96d7a18db487b955505da7fa674b(env, config: dict):
    """Get slide dimensions to determine orientation."""
    import tempfile
    import os
    from pptx import Presentation
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        width = prs.slide_width
        height = prs.slide_height
        is_portrait = height > width
        return {'width': int(width), 'height': int(height), 'is_portrait': is_portrait}
    finally:
        os.unlink(tmp_path)

def get_pptx_italic_check__979792a60bf386d8e1c5a4d702d19829(env, config: dict):
    """Check if text shapes on specified slides have italic formatting."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/43_1.pptx')
    slide_indices = config.get('slide_indices', [0, 1])
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        results = {}
        for idx in slide_indices:
            if idx >= len(prs.slides):
                results[f'slide_{idx}'] = {'error': 'Slide not found'}
                continue
            slide = prs.slides[idx]
            italic_runs = 0
            total_runs = 0
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text.strip():
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                total_runs += 1
                                if run.font.italic:
                                    italic_runs += 1
            results[f'slide_{idx}'] = {'italic_runs': italic_runs, 'total_runs': total_runs}
        return results
    finally:
        os.unlink(tmp_path)

def get_pptx_text_alignment__a27e67f69ab078486669e5f21f103d16(env, config: dict):
    """Get paragraph alignment of a textbox matching text_pattern on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        text_pattern = config.get('text_pattern', '')
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        align_map = {PP_ALIGN.LEFT: 'LEFT', PP_ALIGN.CENTER: 'CENTER', PP_ALIGN.RIGHT: 'RIGHT', PP_ALIGN.JUSTIFY: 'JUSTIFY'}
        for shape in slide.shapes:
            if shape.has_text_frame and text_pattern.lower() in shape.text_frame.text.lower():
                alignments = []
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        align_val = para.alignment
                        align_str = align_map.get(align_val, 'NONE')
                        alignments.append(align_str)
                return {'alignments': alignments, 'shape_name': shape.name, 'text': shape.text_frame.text[:100]}
        return {'error': f'No shape with text "{text_pattern}" on slide {slide_index + 1}'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_hidden_status__266b91e28f8289ba603f3990e8aa8eb6(env, config: dict):
    """Get hidden status of a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Forests.pptx')
    slide_index = config.get('slide_index', 3)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        if slide_index >= slide_count:
            return {'error': f'Slide index {slide_index} out of range (total {slide_count})'}
        slide = prs.slides[slide_index]
        show_attr = slide._element.get('show')
        is_hidden = show_attr == '0'
        title = ''
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        title = text
                        break
            if title:
                break
        return {'slide_count': slide_count, 'slide_index': slide_index, 'is_hidden': is_hidden, 'slide_title': title}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_picture_count__48004e0f145da635f1d1ef72ec0c75ce(env, config: dict):
    """Get the count of picture shapes on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        picture_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
        return {'picture_count': picture_count}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_subtitle__b0063e78116208a0699a2e91287a92ca(env, config: dict):
    """Get the subtitle text (PlaceHolder 2) of a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Writing-Outlines.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        subtitle_text = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 2':
                subtitle_text = shape.text_frame.text.strip()
                break
        return {'subtitle': subtitle_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_pic_height_and_fonts__5b53387f975ab04f84e7061a6b89939a(env, config: dict):
    """Get picture height on a target slide and font sizes of textboxes on another slide."""
    from pptx import Presentation
    from pptx.util import Emu, Cm, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/42_2.pptx')
    pic_slide_index = config.get('pic_slide_index', 4)
    font_slide_index = config.get('font_slide_index', 2)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        pic_height_cm = None
        if pic_slide_index < len(slides):
            slide = slides[pic_slide_index]
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic_height_cm = round(shape.height / 360000, 2)
                    break
        font_sizes = []
        if font_slide_index < len(slides):
            slide = slides[font_slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size is not None:
                                font_sizes.append(round(run.font.size / 12700, 1))
        return {'pic_height_cm': pic_height_cm, 'font_sizes': font_sizes}
    finally:
        os.unlink(tmp_path)

def get_impress_title_size_content_bold_bg__80858354e7972b18d1e56bf385980263(env, config: dict):
    """Get title font size, content bold status, and background color from pptx."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/16_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        title_shape = slide.shapes[0]
        title_size_pt = None
        if title_shape.has_text_frame:
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        title_size_pt = run.font.size / 12700
                    break
                break
        content_shape = slide.shapes[1]
        content_bold = None
        if content_shape.has_text_frame:
            for para in content_shape.text_frame.paragraphs:
                for run in para.runs:
                    content_bold = run.font.bold
                    break
                break
        bg_color = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_color = str(fill.fore_color.rgb)
            except Exception:
                bg_color = None
        return {'title_size_pt': title_size_pt, 'content_bold': content_bold, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_impress_title_font_color__dcc119efe7774eb3168a0af998c68532(env, config: dict):
    """Get the font color of the title text on a specified slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/71_6.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            return {'color_rgb': str(run.font.color.rgb)}
                return {'error': 'No runs with color found in title'}
        return {'error': 'Title placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_title__94c866527f562343eb1ed6d961346913(env, config: dict):
    """Get title text and alignment from a specific slide in a PPTX file."""
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/22_6.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                text = shape.text_frame.text.strip()
                alignment = None
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.alignment == PP_ALIGN.CENTER:
                        alignment = 'center'
                    elif para.alignment == PP_ALIGN.RIGHT:
                        alignment = 'right'
                    elif para.alignment == PP_ALIGN.LEFT:
                        alignment = 'left'
                    else:
                        alignment = 'left'
                return {'text': text, 'alignment': alignment}
        return {'error': 'Title placeholder not found'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_count_and_titles__b3862aa49d8968cf4ed21032ba8b229f(env, config: dict):
    """Get slide count and all slide titles from a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/Forests.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        titles = []
        for slide in prs.slides:
            title = ''
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            title = text
                            break
                if title:
                    break
            titles.append(title)
        return {'slide_count': slide_count, 'titles': titles}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_count_and_titles__3a9a25804cb58719a25a23e52e7e124a(env, config: dict):
    """Get slide count and title text for each slide."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_info = []
        for slide in prs.slides:
            title_text = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        title_text = text
                        break
            slides_info.append({'title': title_text})
        return {'slide_count': len(prs.slides), 'slides': slides_info}
    finally:
        os.unlink(tmp_path)

def get_pptx_picture_position__77f9e57f0af70c9666efd508929f72c3(env, config: dict):
    """Get picture position on a specific slide."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/43_1.pptx')
    slide_index = config.get('slide_index', 1)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': 'Slide not found'}
        slide = prs.slides[slide_index]
        slide_height = prs.slide_height
        slide_width = prs.slide_width
        pictures = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append({'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'bottom': shape.top + shape.height})
        return {'pictures': pictures, 'slide_height': slide_height, 'slide_width': slide_width, 'picture_count': len(pictures)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__6c6be158040f1b190133813bff96d3a4(env, config: dict):
    """Get title text from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = ''
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title_text = text
                    break
        return {'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_count__f16a19b8aecc56a3153baf7aa9668d95(env, config: dict):
    """Get the number of slides in a pptx file."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        return {'slide_count': slide_count}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_title_text__1b7b08a28b40e394d9044a5b56ecfe62(env, config: dict):
    """Get the title text from a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        if slide.shapes.title and slide.shapes.title.text:
            return {'title_text': slide.shapes.title.text.strip()}
        for shape in slide.placeholders:
            if shape.text and shape.text.strip():
                return {'title_text': shape.text.strip()}
        return {'title_text': ''}
    finally:
        os.unlink(tmp_path)

def get_pptx_image_props__19c7a30018144cafc9f4ebd4b22910ef(env, config: dict):
    """Get image properties (count, width, height) from a specific slide."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_path = config.get('path', '/home/user/Desktop/31_2.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        images = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({'width_emu': shape.width, 'height_emu': shape.height, 'width_cm': round(shape.width / 360000, 2), 'height_cm': round(shape.height / 360000, 2), 'name': shape.name})
        return {'image_count': len(images), 'images': images}
    finally:
        os.unlink(tmp_path)

def get_slide1_badge_exists__6b4bf280e2d30f49e851d7a1a643ef02(env, config: dict):
    """Check if a square badge image exists on the right side of slide 1."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        badge_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                try:
                    w_cm = shape.width / 360000
                    h_cm = shape.height / 360000
                    x_cm = shape.left / 360000
                    if abs(w_cm - h_cm) < 1.0 and 5.0 < w_cm < 9.0 and (x_cm > 12.0):
                        badge_found = True
                        break
                except Exception:
                    pass
        return {'badge_exists': badge_found}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_italic__39e2309fd148a4ecc88950c3fc66ef3b(env, config: dict):
    """Check if text in a specific placeholder has italic formatting."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/154_3.pptx')
    shape_name = config.get('shape_name', 'PlaceHolder 2')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                total_runs = 0
                italic_runs = 0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            total_runs += 1
                            if run.font.italic:
                                italic_runs += 1
                return {'total_runs': total_runs, 'italic_runs': italic_runs, 'all_italic': total_runs > 0 and italic_runs == total_runs}
        return {'error': f"Shape '{shape_name}' not found"}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_title_font__05b860587b204cf23347a8d6b2a21320(env, config: dict):
    """Get title text and font name from a specific slide in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/9_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        font_name = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type is not None and int(ph_type) == 1:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                title_text = para.text.strip()
                                for run in para.runs:
                                    if run.font.name is not None:
                                        font_name = run.font.name
                                break
        return {'title_text': title_text, 'font_name': font_name}
    finally:
        os.unlink(tmp_path)

def get_slide1_image_width__d195f7f7ed997b32944186b135a8c2d3(env, config: dict):
    """Get the width of images on slide 1 to check if any has been resized to target width."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        images_info = []
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                try:
                    w_cm = shape.width / 360000
                    h_cm = shape.height / 360000
                    x_cm = shape.left / 360000
                    y_cm = shape.top / 360000
                    images_info.append({'name': shape.name, 'width_cm': round(w_cm, 2), 'height_cm': round(h_cm, 2), 'left_cm': round(x_cm, 2), 'top_cm': round(y_cm, 2)})
                except Exception:
                    pass
        return {'images': images_info}
    finally:
        os.unlink(tmp_path)

def get_slide_title_font__7f1ba412564f3dc4322d685ce4a3274f(env, config: dict):
    """Get the font names used in the title placeholder on a specific slide."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_idx = config.get('slide_index', 0)
        if slide_idx >= len(prs.slides):
            return {'error': f'Slide index {slide_idx} out of range'}
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                if shape.placeholder_format.idx in (0, 15):
                    font_names = []
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                font_names.append(run.font.name)
                    return {'title_text': shape.text_frame.text, 'font_names': font_names}
        return {'error': 'No title placeholder found on slide'}
    finally:
        os.unlink(tmp_path)

def get_impress_title_font_size__3af9d622838a82de0731abee0010dddc(env, config: dict):
    """Get title text and font size from a specific slide in a PPTX file."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/9_1.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = None
        font_size_pt = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    if ph_type is not None and int(ph_type) == 1:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                title_text = para.text.strip()
                                for run in para.runs:
                                    if run.font.size is not None:
                                        font_size_pt = run.font.size / Pt(1)
                                break
        return {'title_text': title_text, 'font_size_pt': font_size_pt}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide5_title_format__e80588656c5c992160d4d54a69813615(env, config: dict):
    """Get font formatting of the title text on slide 5."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/39_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[4]
        for shape in slide.shapes:
            if shape.has_text_frame and 'Infographic' in shape.text:
                results = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        size_pt = None
                        if run.font.size:
                            size_pt = run.font.size.pt
                        results.append({'text': run.text, 'bold': run.font.bold, 'size_pt': size_pt})
                return {'runs': results}
        return {'error': 'Title text not found on slide 5'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title_with_color__fd2e510c274f3f89d65383d731b18828(env, config: dict):
    """Get title text and font color from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = ''
        font_color = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title_text = text
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb:
                                font_color = str(run.font.color.rgb)
                            break
                        break
                    break
        return {'title': title_text, 'font_color': font_color}
    finally:
        os.unlink(tmp_path)

def get_impress_italic_subtitle__910b283898dbff47da6b9e5280c358bb(env, config: dict):
    """Get title italic status and subtitle text from a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        title_italic = None
        subtitle_text = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if 'Should You Consider Studying Abroad' in text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                title_italic = run.font.italic
                                break
                        if title_italic is not None:
                            break
                elif text and subtitle_text is None:
                    subtitle_text = text
        return {'title_italic': title_italic, 'subtitle_text': subtitle_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_bg_color__6efa38f597117d6ec1022844076cbdb0(env, config: dict):
    """Get the background fill color of a specific slide in a PPTX file."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                color = fill.fore_color
                if color and color.rgb:
                    rgb = color.rgb
                    return {'fill_type': str(fill.type), 'color_rgb': str(rgb), 'r': rgb[0] if hasattr(rgb, '__getitem__') else int(str(rgb)[:2], 16), 'g': rgb[1] if hasattr(rgb, '__getitem__') else int(str(rgb)[2:4], 16), 'b': rgb[2] if hasattr(rgb, '__getitem__') else int(str(rgb)[4:6], 16)}
            except Exception:
                pass
        return {'fill_type': str(fill.type), 'color_rgb': None}
    finally:
        os.unlink(tmp_path)

def get_slide_has_bg_image__54b7105f00ae5aefd1a4f4e9a1686c3b(env, config: dict):
    """Check if a specific slide has a background image set."""
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    try:
        file_bytes = env.controller.get_file(ppt_file_path)
        if not file_bytes:
            return {'has_bg_image': False, 'error': 'File not found'}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as myzip:
                slide_xml = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
                if slide_xml not in myzip.namelist():
                    return {'has_bg_image': False, 'error': 'Slide not found'}
                with myzip.open(slide_xml) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                    image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                    attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                    for child in root.iter(bg_tag):
                        for element in child.iter(image_tag):
                            image_id = element.attrib.get(attr_tag, '')
                            if image_id:
                                return {'has_bg_image': True, 'image_id': image_id}
                return {'has_bg_image': False}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error('Error checking slide background: %s', str(e))
        return {'has_bg_image': False, 'error': str(e)}

def get_pptx_slide_title__24e166a4a3628d77570c2919341fa3d0(env, config: dict):
    """Get the title text from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == 'PlaceHolder 1' and shape.has_text_frame:
                title_text = shape.text_frame.text.strip()
                return {'title_text': title_text}
        return {'title_text': ''}
    finally:
        os.unlink(tmp_path)

def get_impress_strikethrough__941b469f6d56a3f2848b88f32a6979cb(env, config: dict):
    """Get strikethrough status of all paragraphs on a specific slide."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        slide = prs.slides[slide_index]
        paragraphs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    all_strike = True
                    has_runs = False
                    for run in para.runs:
                        if run.text.strip():
                            has_runs = True
                            if not run.font.strikethrough:
                                all_strike = False
                                break
                    paragraphs.append({'text': text, 'strikethrough': all_strike and has_runs})
        return {'paragraphs': paragraphs}
    finally:
        os.unlink(tmp_path)

def get_pptx_transitions__a8621b66c753635f21d8c612a865d275(env, config: dict):
    """Get transition types for all slides in a PPTX file."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.oxml.ns import qn
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        transitions = []
        for slide in prs.slides:
            trans_elem = slide._element.find(qn('p:transition'))
            if trans_elem is not None:
                found = False
                for child in trans_elem:
                    tag = child.tag
                    if '}' in tag:
                        tag = tag.split('}')[1]
                    transitions.append(tag.lower())
                    found = True
                    break
                if not found:
                    transitions.append(None)
            else:
                transitions.append(None)
        return {'transitions': transitions, 'total_slides': len(prs.slides)}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_underline__935b1aaa61bb5f50d9e66c69d9b5829b(env, config: dict):
    """Check if text in a specific placeholder has underline formatting."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/154_3.pptx')
    shape_name = config.get('shape_name', 'PlaceHolder 2')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                total_runs = 0
                underlined_runs = 0
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            total_runs += 1
                            if run.font.underline:
                                underlined_runs += 1
                return {'total_runs': total_runs, 'underlined_runs': underlined_runs, 'all_underlined': total_runs > 0 and underlined_runs == total_runs}
        return {'error': f"Shape '{shape_name}' not found"}
    finally:
        os.unlink(tmp_path)

def get_pptx_shape_dimensions__dc32e28ff1403085a7367144a8c3b449(env, config: dict):
    """Get dimensions of specified shapes from a PPTX file."""
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_config = config.get('slides', [])
        results = {}
        for sc in slides_config:
            slide_idx = sc['slide_index']
            shape_name = sc['shape_name']
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.name == shape_name:
                    height_cm = round(shape.height / 360000, 2)
                    width_cm = round(shape.width / 360000, 2)
                    results[f'slide_{slide_idx + 1}_height'] = height_cm
                    results[f'slide_{slide_idx + 1}_width'] = width_cm
                    break
        return results
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title__f01c4fcc320219e33dd1b3dfc752a5b6(env, config: dict):
    """Get title text from a specific slide in a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        title_text = ''
        for shape in slide.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    title_text = text
                    break
        return {'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_orientation__f921e5b58a724c7baa4b59415c1c4019(env, config: dict):
    """Get slide orientation (portrait vs landscape) from a pptx file."""
    import tempfile
    import os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        width = int(prs.slide_width)
        height = int(prs.slide_height)
        is_portrait = height > width
        return {'width': width, 'height': height, 'is_portrait': is_portrait}
    finally:
        os.unlink(tmp_path)

def get_impress_title_and_bg__ac120f719f06d988c1c09c7022ee4b50(env, config: dict):
    """Get slide 5 title text and slide 6 background color from the presentation."""
    from pptx import Presentation
    from pptx.util import Pt
    file_path = config.get('path', '/home/user/Desktop/13_0.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        slide5_title = None
        if len(slides) >= 5:
            slide5 = slides[4]
            for shape in slide5.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text and shape.text_frame.paragraphs[0].runs:
                            run = shape.text_frame.paragraphs[0].runs[0]
                            if run.font.size and run.font.size >= 800000:
                                slide5_title = text
                                break
                if slide5_title:
                    break
        slide6_bg = None
        if len(slides) >= 6:
            slide6 = slides[5]
            fill = slide6.background.fill
            if fill.type is not None:
                try:
                    slide6_bg = str(fill.fore_color.rgb)
                except Exception:
                    slide6_bg = None
        return {'slide5_title': slide5_title, 'slide6_bg_color': slide6_bg}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_count_and_title__e6d2505b7a2f1323c83204293d9b2ff3(env, config: dict):
    """Get the slide count and the title of the last slide from a pptx file."""
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        last_slide_title = ''
        if slide_count > 0:
            last_slide = prs.slides[-1]
            for shape in last_slide.shapes:
                if shape.has_text_frame and 'PlaceHolder 1' in shape.name:
                    last_slide_title = shape.text.strip()
                    break
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0:
                        last_slide_title = shape.text.strip() if shape.has_text_frame else ''
                        break
            if not last_slide_title:
                for shape in last_slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        last_slide_title = shape.text.strip()
                        break
        return {'slide_count': slide_count, 'last_slide_title': last_slide_title}
    finally:
        os.unlink(tmp_path)

def get_impress_strikethrough__0be9a8c316c2b4d9a9ba4471962b751f(env, config: dict):
    """Get strikethrough status of all paragraphs on a specific slide."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        slide = prs.slides[slide_index]
        paragraphs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    all_strike = True
                    has_runs = False
                    for run in para.runs:
                        if run.text.strip():
                            has_runs = True
                            if not run.font.strikethrough:
                                all_strike = False
                                break
                    paragraphs.append({'text': text, 'strikethrough': all_strike and has_runs})
        return {'paragraphs': paragraphs}
    finally:
        os.unlink(tmp_path)

def get_pptx_all_slides_bg_color__88a6aaaf9999692a07033829e3809d82(env, config: dict):
    """Get the background fill color of all slides in a PPTX file."""
    from pptx import Presentation
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides_colors = []
        for (i, slide) in enumerate(prs.slides):
            bg = slide.background
            fill = bg.fill
            color_info = {'slide_index': i, 'color_rgb': None}
            if fill.type is not None:
                try:
                    color = fill.fore_color
                    if color and color.rgb:
                        color_info['color_rgb'] = str(color.rgb).upper()
                except Exception:
                    pass
            slides_colors.append(color_info)
        return {'slides': slides_colors, 'total_slides': len(prs.slides)}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_state__e5c61d88a7cca7248552b922fe7c3602_qw35sft2_3635603e(env, config: dict):
    """Get table state from a specific slide in the PPTX file."""
    import tempfile, os
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    slide_index = config.get('slide_index', 2)
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range (total: {len(prs.slides)})'}
            slide = prs.slides[slide_index]
            tables = []
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    first_cell_text = ''
                    if rows > 0 and cols > 0:
                        first_cell_text = tbl.cell(0, 0).text_frame.text.strip()
                    tables.append({'rows': rows, 'cols': cols, 'first_cell_text': first_cell_text})
            return {'slide_index': slide_index, 'table_count': len(tables), 'tables': tables}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_count__d96a0661e83ae9239b523060c49e36f2_qw35sft2_982cb0b0(env, config: dict):
    """Get the number of slides in a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/AM_Last_Page_Template.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            return {'slide_count': len(prs.slides)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_table_row0__27882b2f85f9d8345d5e8153a8f07379_qw35sft2_49026ab3(env, config: dict):
    """Get the first row of a table on a specific slide from a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/33_1.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 3)
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.shape_type == 19:
                    row0 = shape.table.rows[0]
                    table_row0 = [cell.text.strip() for cell in row0.cells]
                    return {'table_row0': table_row0}
            return {'error': 'No table found on specified slide'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_summary_notes__f1eef050d638a5408c254c4e620a5302_qw35sft2_7b56b6c6(env, config: dict):
    """Get slide count and speaker notes text of the last slide in Forests.pptx."""
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/Forests.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        last_slide = prs.slides[-1]
        notes_text = ''
        try:
            if last_slide.has_notes_slide:
                notes_tf = last_slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text.strip()
        except Exception:
            notes_text = ''
        return {'slide_count': slide_count, 'last_slide_notes': notes_text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_first_slide_text__b9dba6723fe45eeec63ed20d4f46d22d_qw35sft2_c385ede6(env, config: dict):
    """Get slide count and text of the first slide from a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/MLA_Workshop_061X_Works_Cited.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_count = len(prs.slides)
            first_texts = []
            if slide_count > 0:
                for shape in prs.slides[0].shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            first_texts.append(t)
            return {'slide_count': slide_count, 'first_slide_text': ' | '.join(first_texts)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_slide2_title_props__dea2c3a0c75be4ec786c2a7bbd9ea59f_qw35sft2_861497b7(env, config: dict):
    """Get title placeholder position and text content on slide 2 of 134_2.pptx."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not installed'}
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide = prs.slides[1]
        title_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title_shape = shape
                break
        if title_shape is None:
            return {'error': 'Title placeholder not found on slide 2'}
        top_cm = title_shape.top / 360000.0
        title_text = title_shape.text_frame.text.strip()
        return {'top_cm': top_cm, 'title_text': title_text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_pptx_slide_transitions__f348aaf139c455284fcbda7cfe577da5_qw35sft2_4f3addf3(env, config: dict):
    """Get the transition type for all slides in the PPTX file."""
    try:
        from pptx import Presentation
        from lxml import etree
        file_path = config.get('path', '/home/user/Desktop/note-taking-strategies.pptx')
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            transitions = []
            nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
            for slide in prs.slides:
                trans_elem = slide._element.find('.//p:transition', nsmap)
                if trans_elem is not None:
                    fade_elem = trans_elem.find('p:fade', nsmap)
                    dissolve_elem = trans_elem.find('p:dissolve', nsmap)
                    if fade_elem is not None:
                        transitions.append('fade')
                    elif dissolve_elem is not None:
                        transitions.append('dissolve')
                    else:
                        children = list(trans_elem)
                        if children:
                            tag = children[0].tag.split('}')[-1] if '}' in children[0].tag else children[0].tag
                            transitions.append(tag)
                        else:
                            transitions.append('none')
                else:
                    transitions.append('none')
            total = len(transitions)
            fade_count = transitions.count('fade')
            return {'transitions': transitions, 'total_slides': total, 'fade_count': fade_count, 'all_fade': fade_count == total and total > 0}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_file_and_transition__5e93f9ddd16cd129b875ad7386e8a5e5_qw35sft2_a0fba5ee(env, config: dict):
    """Get state of pre.pptx on Desktop: file existence and transitions of all slides."""
    import tempfile, os
    from pptx import Presentation
    from pptx.oxml.ns import qn
    desktop_path = '/home/user/Desktop/pre.pptx'
    file_bytes = env.controller.get_file(desktop_path)
    if not file_bytes:
        return {'file_exists': False, 'all_transitions': [], 'slide_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        transitions = []
        for slide in slides:
            transition_el = slide._element.find(qn('p:transition'))
            if transition_el is not None:
                children = list(transition_el)
                if children:
                    tag = children[0].tag.split('}')[-1] if '}' in children[0].tag else children[0].tag
                    transitions.append(tag)
                else:
                    transitions.append('present')
            else:
                transitions.append(None)
        return {'file_exists': True, 'all_transitions': transitions, 'slide_count': len(slides)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_image_info__7a22b4a6c189fa88dd4ad54ac62a04cd_qw35sft2_2b48f3ad(env, config: dict):
    """Get image information (count and sizes) from a specified slide in a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/31_2.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 1)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range, total slides: {len(prs.slides)}'}
            slide = prs.slides[slide_index]
            images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({'name': shape.name, 'width_cm': round(shape.width / 360000, 4), 'height_cm': round(shape.height / 360000, 4), 'left_cm': round(shape.left / 360000, 4), 'top_cm': round(shape.top / 360000, 4)})
            return {'image_count': len(images), 'images': images, 'slide_index': slide_index}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_presenter_console_state__c6ff7e494d134a42d6e5c420e00bf4fd_qw35sft2_0e3dd544(env, config: dict):
    """Check if LibreOffice Impress Presenter Console is disabled via registrymodifications.xcu."""
    xcu_path = '/home/user/.config/libreoffice/4/user/registrymodifications.xcu'
    file_bytes = env.controller.get_file(xcu_path)
    if not file_bytes:
        return {'error': 'XCU file not found', 'presenter_disabled': False}
    try:
        content = file_bytes.decode('utf-8') if isinstance(file_bytes, bytes) else str(file_bytes)
        pattern = 'PresenterScreenEnabled[^<]*<value>(false|true)</value>'
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            presenter_disabled = match.group(1).lower() == 'false'
        else:
            presenter_disabled = False
        return {'presenter_disabled': presenter_disabled}
    except Exception as e:
        return {'error': str(e), 'presenter_disabled': False}

def get_pptx_bg_fill_type__7f279c8dbe68c5780c5b5f60310b951b_qw35sft2_67387c51(env, config: dict):
    """Get the background fill type name of slide 0 in the PPTX file."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        fill = slide.background.fill
        fill_type_name = fill.type.name if fill.type is not None else 'NONE'
        return {'fill_type': fill_type_name}
    finally:
        os.unlink(tmp_path)

def get_pptx_transitions__85331dc0cf4d821fa2bad3e27657d683_qw35sft2_0149f70f(env, config: dict):
    """Get transition information for all slides in a PPTX file.

    Returns a dict with per-slide transition types parsed from OOXML.
    Transition element children map to types like 'fade', 'dissolve', 'blinds', etc.

    Returns:
      {
        'num_slides': int,
        'transitions': [{'slide': int, 'type': str}, ...],
        'num_with_transitions': int
      }
    """
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/lec17-gui-events.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        _PML_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        transitions = []
        for i, slide in enumerate(prs.slides):
            trans_el = slide._element.find(f'{{{_PML_NS}}}transition')
            if trans_el is None:
                transitions.append({'slide': i, 'type': 'none'})
                continue
            children = [c for c in trans_el if not callable(c)]
            if children:
                raw_tag = children[0].tag
                t_type = raw_tag.split('}')[-1] if '}' in raw_tag else raw_tag
            else:
                t_type = 'cut'
            transitions.append({'slide': i, 'type': t_type.lower()})
        num_with = sum((1 for t in transitions if t['type'] not in ('none', 'cut')))
        return {'num_slides': len(prs.slides), 'transitions': transitions, 'num_with_transitions': num_with}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_pptx_slide_bg_color__1a5b2a2b44ee739e3486ddbc20ac1c87_qw35sft2_88959acb(env, config: dict):
    """Get the background solid fill color of a specific slide in the PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.dml import MSO_THEME_COLOR
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/13_0.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 2)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range'}
            slide = prs.slides[slide_index]
            bg = slide.background
            fill = bg.fill
            if fill.type is None:
                return {'error': 'No background fill set', 'color_hex': None}
            from pptx.enum.dml import MSO_FILL
            if str(fill.type) != 'SOLID (1)':
                return {'error': f'Background fill is not solid: {fill.type}', 'color_hex': None}
            try:
                rgb = fill.fore_color.rgb
                hex_str = str(rgb)
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
                return {'color_hex': hex_str, 'r': r, 'g': g, 'b': b, 'slide_index': slide_index}
            except Exception as e:
                return {'error': f'Could not read color: {e}', 'color_hex': None}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_bg_fill__30ec500dc9631258b0ce307c50b98145_qw35sft2_26833ac2(env, config: dict):
    """Get slide background fill type and color from a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/214_9.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_idx = config.get('slide_index', 0)
            if slide_idx >= len(prs.slides):
                return {'error': f'Slide index {slide_idx} out of range'}
            slide = prs.slides[slide_idx]
            fill = slide.background.fill
            fill_type = fill.type
            result = {'fill_type_value': fill_type.real if fill_type is not None else None, 'fill_type_name': fill_type.name if fill_type is not None else None, 'is_solid': False}
            if fill_type is not None and fill_type.real == 1:
                result['is_solid'] = True
                try:
                    rgb = fill.fore_color.rgb
                    hex_str = str(rgb)
                    result['rgb'] = hex_str
                    result['r'] = int(hex_str[0:2], 16)
                    result['g'] = int(hex_str[2:4], 16)
                    result['b'] = int(hex_str[4:6], 16)
                except Exception:
                    pass
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_slide4_contact_shapes__5eae632a29df23abd3f96eee0032e6dc_qw35sft2_c546a846(env, config: dict):
    """
    Get the count and texts of contact info shapes in slide 4 of 21_0.pptx.
    Original slide 4 has personal info: address, phone, email, social handle, website,
    plus icon shapes (FREEFORM type). Returns shape counts and text content.
    """
    import tempfile, os
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {'error': 'python-pptx not installed'}
    file_path = config.get('path', '/home/user/Desktop/21_0.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[3]
        texts = []
        shape_count = len(slide.shapes)
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        all_text = ' '.join(texts)
        has_phone = '123-456-7890' in all_text
        has_address = '123 Anywhere St' in all_text
        has_email = 'hello@reallygreatsite.com' in all_text
        has_social = '@reallygreatsite' in all_text
        has_website = 'reallygreatsite.com' in all_text and '@reallygreatsite' not in all_text.replace('reallygreatsite.com', '')
        return {'shape_count': shape_count, 'texts': texts, 'has_phone': has_phone, 'has_address': has_address, 'has_email': has_email, 'has_social': has_social, 'has_website': 'reallygreatsite.com' in all_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_stretch_and_shapes__bd30ea01dba8f440569b065c94e4f32f_qw35sft2_7b4096a9(env, config: dict):
    import tempfile, os
    from pptx import Presentation
    pptx_path = config.get('path', '/home/user/Desktop/CPD_Background_Investigation_Process.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        slide_width_cm = prs.slide_width / 914400 * 2.54
        slide_height_cm = prs.slide_height / 914400 * 2.54
        best_shape = None
        best_width = 0.0
        for shape in slide.shapes:
            if shape.shape_type == 13 and 'Picture 2' in shape.name:
                w = shape.width / 914400 * 2.54
                if w > best_width:
                    best_width = w
                    best_shape = shape
        image_width_cm = best_shape.width / 914400 * 2.54 if best_shape else 0.0
        image_height_cm = best_shape.height / 914400 * 2.54 if best_shape else 0.0
        image_left_cm = best_shape.left / 914400 * 2.54 if best_shape else -1.0
        image_top_cm = best_shape.top / 914400 * 2.54 if best_shape else -1.0
        picture3_exists = any((shape.name == 'Picture 3' for shape in slide.shapes))
        return {'image_width_cm': round(image_width_cm, 2), 'image_height_cm': round(image_height_cm, 2), 'image_left_cm': round(image_left_cm, 2), 'image_top_cm': round(image_top_cm, 2), 'slide_width_cm': round(slide_width_cm, 2), 'slide_height_cm': round(slide_height_cm, 2), 'picture3_exists': picture3_exists}
    finally:
        os.unlink(tmp_path)

def get_impress_underline_state__fa317ec6034191f7b02a906f78ff4b09_qw35sft2_128f5711(env, config: dict):
    """
    Get the underline state of all text runs in a specific shape on a specific slide.
    config keys:
      path: VM path to the pptx file
      slide_index: 0-based slide index
      shape_name: name of the shape to check
    Returns dict with 'all_underlined' bool and 'run_count' int.
    """
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/43_1.pptx')
    slide_index = config.get('slide_index', 1)
    shape_name = config.get('shape_name', 'Google Shape;573;p15')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        if slide_index >= len(slides):
            return {'error': f'Slide {slide_index} not found'}
        slide = slides[slide_index]
        for shape in slide.shapes:
            if shape.name == shape_name:
                if not hasattr(shape, 'text_frame'):
                    return {'error': 'Shape has no text_frame'}
                runs = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        runs.append(run.font.underline)
                if not runs:
                    return {'error': 'No text runs found', 'run_count': 0}
                all_underlined = all((u is True for u in runs))
                return {'all_underlined': all_underlined, 'run_count': len(runs), 'underline_values': runs}
        return {'error': f'Shape {shape_name!r} not found on slide {slide_index}'}
    finally:
        os.unlink(tmp_path)

def get_impress_audio_slide1_trans_slide3__1b01b098ee4b6dfab8bc1b2d24723a1f_qw35sft2_18b7000f(env, config: dict):
    """Check if slide 1 has audio and if slide 3 has any transition applied."""
    ppt_path = config.get('path', '/home/user/Desktop/Mady_and_Mia_Baseball.pptx')
    file_bytes = env.controller.get_file(ppt_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_audio_slide1': False, 'has_transition_slide3': False}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_audio_slide1 = False
        has_transition_slide3 = False
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            rels_file = 'ppt/slides/_rels/slide1.xml.rels'
            if rels_file in zf.namelist():
                with zf.open(rels_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    for rel in root.findall('r:Relationship', ns):
                        if 'audio' in rel.attrib.get('Type', ''):
                            has_audio_slide1 = True
                            break
            slide3_file = 'ppt/slides/slide3.xml'
            if slide3_file in zf.namelist():
                with zf.open(slide3_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    trans = root.find('.//{%s}transition' % pns)
                    if trans is not None:
                        has_transition_slide3 = True
        return {'has_audio_slide1': has_audio_slide1, 'has_transition_slide3': has_transition_slide3}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_info__da57195bec8b00f5d49d80dfe28b0049_qw35sft2_ccd6c0f7(env, config: dict):
    """Get slide count and text frame info from a PPTX presentation on the VM."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/simple.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': f'File not found: {path}'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        slides_with_text_frames = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slides_with_text_frames += 1
                    break
        return {'slide_count': slide_count, 'slides_with_text_frames': slides_with_text_frames}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_bullet_newpara__12c621423b2418d27e86fd746731530a_qw35sft2_457c566e(env, config: dict):
    """Get bullet status on first para and text of second paragraph in content."""
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/69_4.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        content_shape = slide.shapes[1]
        tf = content_shape.text_frame
        paras = tf.paragraphs
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        first_para_has_bullet = False
        if paras:
            pPr = paras[0]._p.find(f'{{{ns}}}pPr')
            if pPr is not None:
                buChar = pPr.find(f'{{{ns}}}buChar')
                first_para_has_bullet = buChar is not None
        second_para_text = ''
        if len(paras) > 1:
            second_para_text = paras[1].text.strip()
        return {'first_para_has_bullet': first_para_has_bullet, 'second_para_text': second_para_text, 'total_paragraphs': len(paras)}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_alignments__6253f908caa7dec8dc9ccbad6c2f7332_qw35sft2_c050a0a8(env, config: dict):
    """Get text alignment for first textboxes on slides 3, 4, and 5 of 38_1.pptx."""
    import tempfile, os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/38_1.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        results = {}
        slide_searches = [(2, 'HELLO', 'slide3_align'), (3, 'WRITE YOUR', 'slide4_align'), (4, 'WRITE AN ORIGINAL', 'slide5_align')]
        for slide_idx, keyword, key in slide_searches:
            slide = prs.slides[slide_idx]
            align_str = None
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and keyword.lower() in shape.text_frame.text.lower():
                    align = shape.text_frame.paragraphs[0].alignment
                    if align == PP_ALIGN.RIGHT:
                        align_str = 'RIGHT'
                    elif align == PP_ALIGN.CENTER:
                        align_str = 'CENTER'
                    elif align == PP_ALIGN.LEFT:
                        align_str = 'LEFT'
                    else:
                        align_str = 'NONE'
                    break
                if shape.shape_type == 6:
                    for sub in shape.shapes:
                        if hasattr(sub, 'text_frame') and keyword.lower() in sub.text_frame.text.lower():
                            align = sub.text_frame.paragraphs[0].alignment
                            if align == PP_ALIGN.RIGHT:
                                align_str = 'RIGHT'
                            elif align == PP_ALIGN.CENTER:
                                align_str = 'CENTER'
                            elif align == PP_ALIGN.LEFT:
                                align_str = 'LEFT'
                            else:
                                align_str = 'NONE'
                            break
                    if align_str is not None:
                        break
            results[key] = align_str
        return results
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_font_size__4b92c85bf54981a99e31cdf8daa555a4_qw35sft2_86e22296(env, config: dict):
    """Get font name and size (pt) of the last slide title in 24_8.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/24_8.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[-1]
        result = {'font_name': None, 'font_size_pt': None}
        for shape in slide.shapes:
            if shape.has_text_frame and shape.placeholder_format is not None and (shape.placeholder_format.idx == 0):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            font_name = run.font.name
                            font_size = run.font.size
                            if font_name is None:
                                try:
                                    from pptx.oxml.ns import qn
                                    pPr = para._p.find(qn('a:pPr'))
                                    if pPr is not None:
                                        defRPr = pPr.find(qn('a:defRPr'))
                                        if defRPr is not None:
                                            font_name = defRPr.get('lang') or None
                                            latin = defRPr.find(qn('a:latin'))
                                            if latin is not None:
                                                font_name = latin.get('typeface') or font_name
                                except Exception:
                                    pass
                            if font_size is None:
                                try:
                                    from pptx.oxml.ns import qn
                                    from pptx.util import Pt
                                    pPr = para._p.find(qn('a:pPr'))
                                    if pPr is not None:
                                        defRPr = pPr.find(qn('a:defRPr'))
                                        if defRPr is not None:
                                            sz_val = defRPr.get('sz')
                                            if sz_val is not None:
                                                font_size = Pt(int(sz_val) / 100.0)
                                except Exception:
                                    pass
                            result['font_name'] = font_name
                            if font_size is not None:
                                result['font_size_pt'] = round(font_size.pt)
                            return result
        return result
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide3_text_colors__9e2b961f79a8b89da2fe1a7f40c57f65_qw35sft2_4dc70f53(env, config: dict):
    """Get all text run font colors from slide 3 of 1_2.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/1_2.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[2]
        colors = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            if run.font.color and run.font.color.type is not None:
                                colors.append(str(run.font.color.rgb).upper())
                            else:
                                colors.append(None)
        return {'colors': colors, 'count': len(colors)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide3_para_texts__c0cb167c6059fd8f9f81e328b9599144_qw35sft2_81111149(env, config: dict):
    """Read paragraph texts from slide 3's content placeholder in Writing-Outlines.pptx."""
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/Writing-Outlines.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[2]
        shape = slide.shapes[1]
        tf = shape.text_frame
        texts = [para.text for para in tf.paragraphs]
        return {'paragraphs': texts}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_multi_title_color__14718ed4f8dc81749072a781b972112d_qw35sft2_09d2b40d(env, config: dict):
    """
    Get the font color of the topmost text shape (title) for multiple slides in a PPTX file.
    config keys: path (VM path), slide_indices (list of 0-based indices)
    Returns: {'slides': {'1': {'color_hex': ..., 'title_text': ...}, ...}, 'error': None}
    """
    from pptx import Presentation
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/4_1.pptx'))
    if not file_bytes:
        return {'error': 'File not found', 'slides': {}}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_indices = config.get('slide_indices', [1, 2])
        result_slides = {}
        for slide_index in slide_indices:
            if slide_index >= len(prs.slides):
                result_slides[str(slide_index)] = {'error': f'Slide index {slide_index} out of range', 'color_hex': None}
                continue
            slide = prs.slides[slide_index]
            shapes_with_text = [(shape.top if shape.top is not None else float('inf'), shape) for shape in slide.shapes if hasattr(shape, 'text_frame') and shape.text.strip()]
            if not shapes_with_text:
                result_slides[str(slide_index)] = {'error': 'No text shapes found', 'color_hex': None}
                continue
            shapes_with_text.sort(key=lambda x: x[0])
            title_shape = shapes_with_text[0][1]
            colors = []
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    if font.color and font.color.type is not None:
                        try:
                            colors.append(str(font.color.rgb).upper())
                        except Exception:
                            colors.append(None)
                    else:
                        colors.append(None)
            result_slides[str(slide_index)] = {'color_hex': colors[0] if colors else None, 'all_colors': colors, 'title_text': title_shape.text.strip()[:80], 'error': None}
        return {'slides': result_slides, 'error': None}
    finally:
        os.unlink(tmp_path)

def get_slide2_notes_state__dd477a39d59856cb4084c0b8f3f8bce8_qw35sft2_a6f6bdb8(env, config: dict):
    """Get the speaker notes text on slide 2 of 164_3.pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'pptx module not available'}
    file_path = config.get('path', '/home/user/Desktop/164_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide = prs.slides[1]
        notes_text = ''
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf:
                notes_text = tf.text.strip()
        return {'notes_text': notes_text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_slide1_title__992b8b598a921e5736a852e261be2237_qw35sft2_7462f87d(env, config: dict):
    """Get the title text of slide 1 in a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/189_4.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                    text = shape.text_frame.text.strip()
                    return {'title_text': text}
            return {'error': 'Title placeholder not found', 'title_text': ''}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_image_quadrant__f123db274ed42f1e244c71177f056c09_qw35sft2_5f5fc1c9(env, config: dict):
    """Get image position (left, top) relative to slide dimensions on Slide 2 of 201_6.pptx."""
    file_path = config.get('path', '/home/user/Desktop/201_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if shape.shape_type == 13:
                return {'image_left': shape.left, 'image_top': shape.top, 'image_width': shape.width, 'image_height': shape.height, 'slide_width': slide_width, 'slide_height': slide_height}
        return {'error': 'No picture found on Slide 2'}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide2_title_state__1225c3617161e6b819a6b6275a5c9749_qw35sft2_b7f0038d(env, config: dict):
    """Get slide 2 title text and alignment from the PPTX file."""
    import tempfile, os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/22_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide2 = prs.slides[1]
        title_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                title_shape = shape
                break
        if title_shape is None:
            for shape in slide2.shapes:
                if shape.has_text_frame:
                    title_shape = shape
                    break
        if title_shape is None:
            return {'title_text': '', 'title_alignment': None}
        tf = title_shape.text_frame
        title_text = tf.text.strip()
        alignment = None
        if tf.paragraphs:
            alignment = tf.paragraphs[0].alignment
        alignment_name = None
        if alignment is not None:
            alignment_name = alignment.name
        return {'title_text': title_text, 'title_alignment': alignment_name}
    finally:
        os.unlink(tmp_path)

def get_impress_slide3_table_pos__db4d3e235d7b0e304c073a25921cbfa2_qw35sft2_c36fb8ca(env, config: dict):
    """Get the position of the table on slide 3 of the PPTX file."""
    import tempfile
    import os
    pptx_path = config.get('path', '/home/user/Desktop/55_10.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        if len(prs.slides) < 3:
            return {'error': f'Expected at least 3 slides, got {len(prs.slides)}'}
        slide3 = prs.slides[2]
        slide_height = prs.slide_height
        table_shape = None
        for shape in slide3.shapes:
            if shape.shape_type == 19:
                table_shape = shape
                break
        if table_shape is None:
            return {'error': 'No table found on slide 3'}
        return {'table_top': table_shape.top, 'table_height': table_shape.height, 'slide_height': slide_height, 'table_bottom_edge': table_shape.top + table_shape.height}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_underline__42b6fd5c06cd718a123cb3a7892d354e_qw35sft2_fe20d734(env, config: dict):
    """Get underline status of all runs in the title shape of slide 1."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/154_3.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide = prs.slides[0]
            title_shape = slide.shapes[0]
            if not title_shape.has_text_frame:
                return {'error': 'Title shape has no text frame'}
            runs_info = []
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    runs_info.append({'text': run.text, 'underline': run.font.underline})
            all_underlined = bool(runs_info) and all((r['underline'] is True for r in runs_info))
            return {'title_runs': runs_info, 'all_underlined': all_underlined}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_slide14_font_sizes__ec3f2f9f447d657b0668843895ae7804_qw35sft2_1714f192(env, config: dict):
    """Get font sizes of the first two non-empty textboxes on slide 14 of 45_1.pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/45_1.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[13]
        sizes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            size_pt = None
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.strip() and run.font.size is not None:
                        size_pt = round(run.font.size / 12700)
                        break
                if size_pt is not None:
                    break
            sizes.append({'text': tf.text[:50], 'size_pt': size_pt})
            if len(sizes) == 2:
                break
        return {'textbox1_size_pt': sizes[0]['size_pt'] if len(sizes) > 0 else None, 'textbox1_text': sizes[0]['text'] if len(sizes) > 0 else None, 'textbox2_size_pt': sizes[1]['size_pt'] if len(sizes) > 1 else None, 'textbox2_text': sizes[1]['text'] if len(sizes) > 1 else None}
    finally:
        os.unlink(tmp_path)

def get_impress_notes_bg_title__cb759dd1cf89c2aec2802f8543d62400_qw35sft2_73e00eef(env, config: dict):
    """Get slide notes, background RGB, and title text from a PPTX."""
    import tempfile
    import os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/181_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        bg_rgb = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_rgb = str(fill.fore_color.rgb)
            except Exception:
                bg_rgb = None
        title_text = ''
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_text = title_shape.text.strip()
        return {'notes_text': notes_text, 'bg_rgb': bg_rgb, 'title_text': title_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_color__39fe73dd35955e4e54b2d721092898e1_qw35sft2_47e52091(env, config: dict):
    """Get the font color of a specific text on a slide in the PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/saa-format-guide.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            target_text = config.get('target_text', '')
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.replace('\n', ' ').strip()
                if target_text and target_text.upper() not in text.upper():
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.type is not None:
                            try:
                                rgb = run.font.color.rgb
                                hex_str = str(rgb)
                                return {'text': text, 'color_hex': hex_str.upper(), 'r': int(hex_str[0:2], 16), 'g': int(hex_str[2:4], 16), 'b': int(hex_str[4:6], 16)}
                            except Exception:
                                pass
            return {'error': f'Text "{target_text}" not found or has no explicit color set'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_single_text_color__c830024b32b30956e96002af134f6f4a_qw35sft2_25f7b637(env, config: dict):
    """Get the font color of a specific textbox in slide 1 of a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/45_2.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            target_text = config.get('target_text', 'LAUNCH')
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.replace('\n', ' ').strip()
                    if target_text.upper() in text.upper():
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.type is not None:
                                    try:
                                        rgb = run.font.color.rgb
                                        hex_str = str(rgb)
                                        r = int(hex_str[0:2], 16)
                                        g = int(hex_str[2:4], 16)
                                        b = int(hex_str[4:6], 16)
                                        return {'text': text, 'color_hex': hex_str, 'r': r, 'g': g, 'b': b}
                                    except Exception:
                                        pass
            return {'error': f'Text "{target_text}" not found or has no explicit color'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_slide2_bg_subtitle_color__2d907c4b088ee7e1bdf99fbd932517ff_qw35sft2_3d0a0bfd(env, config: dict):
    """Get slide 2 background color and subtitle text color from 71_6.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/71_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[1]
        bg_color = None
        try:
            fill = slide.background.fill
            if fill.type is not None and int(fill.type) == 1:
                bg_color = str(fill.fore_color.rgb).upper()
        except Exception:
            pass
        subtitle_color = None
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        if len(text_shapes) >= 2:
            body_shape = text_shapes[1]
            for para in body_shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        subtitle_color = str(run.font.color.rgb).upper()
                    except Exception:
                        pass
                    if subtitle_color:
                        break
                if subtitle_color:
                    break
        return {'bg_color': bg_color, 'subtitle_color': subtitle_color}
    finally:
        os.unlink(tmp_path)

def get_pptx_title_subtitle_font__761afd0460bdb180560f4c7116933937_qw35sft2_0b2c4583(env, config: dict):
    """Get title text, font name, and subtitle text from slide 1 of 9_1.pptx."""
    try:
        from pptx import Presentation
        from pptx.util import PP_PLACEHOLDER
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/9_1.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            slide = prs.slides[slide_index]
            title_text = None
            font_name = None
            subtitle_text = None
            shapes_by_idx = {}
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    idx = shape.placeholder_format.idx
                    shapes_by_idx[idx] = shape
            if 0 in shapes_by_idx:
                title_shape = shapes_by_idx[0]
                title_text = title_shape.text_frame.text.strip()
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            font_name = run.font.name
                        break
                    if font_name:
                        break
            if 1 in shapes_by_idx:
                sub_shape = shapes_by_idx[1]
                subtitle_text = sub_shape.text_frame.text.strip()
            if title_text is None:
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            texts.append((shape, txt))
                if texts:
                    title_shape, title_text = texts[0]
                    for para in title_shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                font_name = run.font.name
                            break
                        if font_name:
                            break
                if len(texts) > 1:
                    _, subtitle_text = texts[1]
            return {'title_text': title_text, 'font_name': font_name, 'subtitle_text': subtitle_text}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_pptx_props__ed92c6da7770c48f6ebca23c6070cbf5_qw35sft2_da03d5c6(env, config: dict):
    """Get slide 3 Group 6 height and slide 6 textbox font sizes from 42_2.pptx."""
    import tempfile, os
    from pptx import Presentation
    from pptx.oxml.ns import qn
    file_bytes = env.controller.get_file('/home/user/Desktop/42_2.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        result = {}
        slide3 = prs.slides[2]
        result['slide3_group6_height_cm'] = None
        for shape in slide3.shapes:
            if shape.name == 'Group 6':
                result['slide3_group6_height_cm'] = round(shape.height / 360000, 4)
                break
        slide6 = prs.slides[5]
        slide6_font_sizes = {}
        for shape in slide6.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text.strip():
                continue
            pts = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        pts.append(int(round(run.font.size.pt)))
                    else:
                        defRPr = para._p.find(qn('a:defRPr'))
                        if defRPr is not None:
                            sz = defRPr.get('sz')
                            if sz:
                                pts.append(int(sz) // 100)
                if not para.runs:
                    defRPr = para._p.find(qn('a:defRPr'))
                    if defRPr is not None:
                        sz = defRPr.get('sz')
                        if sz:
                            pts.append(int(sz) // 100)
            if pts:
                slide6_font_sizes[shape.name] = pts
        result['slide6_font_sizes'] = slide6_font_sizes
        return result
    finally:
        os.unlink(tmp_path)

def get_pptx_content_and_table__7f43e77e7e3ce83ed182df2ac03d3d3a_qw35sft2_1259b210(env, config: dict):
    """Get Slide 2 content placeholder text and the first cell (row 0, col 0) of the table on Slide 2."""
    import tempfile
    import os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/109_4.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide2_content = ''
        table_cell_00 = ''
        if len(prs.slides) >= 2:
            slide2 = prs.slides[1]
            for shape in slide2.shapes:
                if shape.has_text_frame and shape.name == 'PlaceHolder 2':
                    slide2_content = shape.text_frame.text.strip()
                elif shape.shape_type == 19:
                    table = shape.table
                    if len(table.rows) > 0 and len(table.rows[0].cells) > 0:
                        table_cell_00 = table.rows[0].cells[0].text.strip()
        return {'slide2_content': slide2_content, 'table_cell_00': table_cell_00}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_picture_size_slide6__f55d52551229ea682c89b1a15474ccff_qw35sft2_cf035c45(env, config: dict):
    """Get picture heights on slides 3, 4, 6 and width on slide 6 from 30_1.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/30_1.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        emu_to_cm = lambda emu: round(emu / 360000.0, 4)
        slide3 = prs.slides[2]
        slide3_target = None
        for shape in slide3.shapes:
            if shape.shape_type == 5:
                if slide3_target is None or shape.height > slide3_target.height:
                    slide3_target = shape
        slide4 = prs.slides[3]
        slide4_target = None
        for shape in slide4.shapes:
            if shape.shape_type == 13 and shape.height > 5 * 360000:
                if slide4_target is None or shape.height > slide4_target.height:
                    slide4_target = shape
        slide6 = prs.slides[5]
        slide6_target = None
        for shape in slide6.shapes:
            if shape.shape_type == 13:
                if slide6_target is None or shape.width > slide6_target.width:
                    slide6_target = shape
        return {'slide3_height_cm': emu_to_cm(slide3_target.height) if slide3_target else None, 'slide4_height_cm': emu_to_cm(slide4_target.height) if slide4_target else None, 'slide6_height_cm': emu_to_cm(slide6_target.height) if slide6_target else None, 'slide6_width_cm': emu_to_cm(slide6_target.width) if slide6_target else None}
    finally:
        os.unlink(tmp_path)

def get_impress_panel_slide_count__ddf9714ddb153a97bbdfb4350e733817_qw35sft2_1a390bcc(env, config: dict):
    """
    Get slide pane visibility and number of slides visible in the Slides panel.

    LibreOffice Impress accessibility tree lists each slide thumbnail with a
    label containing its index (e.g. 'Slide 1', 'Slide 2').  We count those
    occurrences to determine the slide count.
    """
    tree = env.controller.get_accessibility_tree()
    if not isinstance(tree, str):
        return {'panel_visible': False, 'slide_count': 0}
    panel_visible = 'Slides' in tree
    slide_entries = re.findall('\\bSlide\\s+\\d+\\b', tree)
    slide_count = len(slide_entries)
    if panel_visible and slide_count == 0:
        slide_count = 1
    return {'panel_visible': panel_visible, 'slide_count': slide_count}

def get_pptx_text_font_size__f00d8a2d9828be4b94677a9606cd7f47_qw35sft2_37d00a00(env, config: dict):
    """Get font sizes of all runs in a specific shape of the first slide."""
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/16_2.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        shape_index = config.get('shape_index', 0)
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
        return {'font_sizes': sizes, 'shape_index': shape_index}
    finally:
        os.unlink(tmp_path)

def get_pptx_orientation__2491631e79810e0b815ac32866fd3274_qw35sft2_0f4fbcfa(env, config: dict):
    """Get slide orientation (width/height) from a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/AM_Last_Page_Template.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            width = prs.slide_width
            height = prs.slide_height
            return {'width_emu': int(width), 'height_emu': int(height), 'is_portrait': int(height) > int(width)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_orientation_state__f5c58fc75d059d716d7a9ec9e8e3967d_qw35sft2_e2628c21(env, config: dict):
    """Get slide count and orientation (portrait or landscape) of Forests.pptx."""
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/Forests.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        is_portrait = prs.slide_height > prs.slide_width
        return {'slide_count': slide_count, 'is_portrait': is_portrait}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_state__99738e8866681b9fdc97e255c87e124d_qw35sft2_2395336c(env, config: dict):
    """Get table state from a specific slide in the PPTX file."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    slide_index = config.get('slide_index', 2)
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range (total: {len(prs.slides)})'}
            slide = prs.slides[slide_index]
            tables = []
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    rows = len(tbl.rows)
                    cols = len(tbl.columns)
                    tables.append({'rows': rows, 'cols': cols})
            return {'slide_index': slide_index, 'table_count': len(tables), 'tables': tables}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_slide2_title_top__ff47db79abc78d9904bd7e17dbb1e3d3_qw35sft2_13ea5fa8(env, config: dict):
    """Get the vertical position (top) of the title placeholder on slide 2 of 134_2.pptx."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not installed'}
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide = prs.slides[1]
        title_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title_shape = shape
                break
        if title_shape is None:
            return {'error': 'Title placeholder not found on slide 2'}
        top_cm = title_shape.top / 360000.0
        height_cm = title_shape.height / 360000.0
        slide_height_cm = prs.slide_height / 360000.0
        return {'top_cm': top_cm, 'height_cm': height_cm, 'slide_height_cm': slide_height_cm}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_pptx_notes_nonempty__d897409e855e2cbbb791ebde03eeef8a_qw35sft2_37fdad70(env, config: dict):
    """Get whether slide 0 notes are non-empty, and return the notes text."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        if not slide.has_notes_slide:
            return {'notes': '', 'nonempty': False}
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        return {'notes': notes_text, 'nonempty': bool(notes_text)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_title_text__c928e38a883d919c295367c551657155_qw35sft2_7b8b2632(env, config: dict):
    """Get the title placeholder text from a specific slide in a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/214_9.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_idx = config.get('slide_index', 1)
            if slide_idx >= len(prs.slides):
                return {'error': f'Slide index {slide_idx} out of range'}
            slide = prs.slides[slide_idx]
            title_text = ''
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    pf = shape.placeholder_format
                    if pf is not None and pf.type.real == 1:
                        title_text = shape.text_frame.text.strip()
                        break
                except Exception:
                    continue
            return {'title_text': title_text}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_image_and_transition__4195ceef20cb2389eadfdf0daf18e8e7_qw35sft2_77e2d736(env, config: dict):
    """Get image info and transition presence from a specified slide in a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/31_2.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 1)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range, total slides: {len(prs.slides)}'}
            slide = prs.slides[slide_index]
            images = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({'name': shape.name, 'width_cm': round(shape.width / 360000, 4), 'height_cm': round(shape.height / 360000, 4)})
            transition_elem = slide._element.find(qn('p:transition'))
            has_transition = transition_elem is not None
            transition_type = None
            if has_transition:
                for ttype in ['fade', 'push', 'wipe', 'split', 'reveal', 'circle', 'diamond', 'dissolve', 'zoom', 'flash']:
                    if transition_elem.find(qn(f'p:{ttype}')) is not None:
                        transition_type = ttype
                        break
                if transition_type is None:
                    transition_type = 'unknown'
            return {'image_count': len(images), 'images': images, 'has_transition': has_transition, 'transition_type': transition_type}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_bg_color__b309be1bc27dcbe28bf5749a28cc9b6f_qw35sft2_f940a712(env, config: dict):
    """Get background color of a specific slide in a PPTX file.

    Returns dict with fill_type, rgb (hex string, 6 chars), and fallback
    master_rgb when the slide itself does not have an explicit solid fill.
    """
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/lec17-gui-events.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'rgb': None, 'master_rgb': None}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range (has {len(prs.slides)} slides)', 'rgb': None, 'master_rgb': None}
        slide = prs.slides[slide_index]
        bg = slide.background
        fill = bg.fill
        result = {'fill_type': fill.type.name if fill.type is not None else 'NONE', 'rgb': None, 'master_rgb': None, 'slide_index': slide_index, 'num_slides': len(prs.slides)}
        if fill.type is not None and fill.type.name == 'SOLID':
            try:
                rgb = fill.fore_color.rgb
                result['rgb'] = str(rgb).upper()
            except Exception as e:
                result['rgb_error'] = str(e)
        if result['rgb'] is None and len(prs.slide_masters) > 0:
            master_fill = prs.slide_masters[0].background.fill
            result['master_fill_type'] = master_fill.type.name if master_fill.type is not None else 'NONE'
            if master_fill.type is not None and master_fill.type.name == 'SOLID':
                try:
                    rgb = master_fill.fore_color.rgb
                    result['master_rgb'] = str(rgb).upper()
                except Exception as e:
                    result['master_rgb_error'] = str(e)
        return result
    except Exception as e:
        return {'error': str(e), 'rgb': None, 'master_rgb': None}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_pptx_slide_subtitle_text__4d840f044cdf667c54faf806c6af73e1_qw35sft2_3db23da2(env, config: dict):
    """Get the subtitle (second text shape) text from a specific slide in the PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/13_0.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 1)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range'}
            slide = prs.slides[slide_index]
            text_shapes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        text_shapes.append(text)
            if len(text_shapes) < 2:
                return {'error': 'Less than 2 text shapes found', 'all_texts': text_shapes}
            return {'subtitle_text': text_shapes[1], 'all_texts': text_shapes}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_last_two_duplicated__ae132e5d53d956fce4051acdf425fcba_qw35sft2_499adc83(env, config: dict):
    """Get slide count and texts of slides 25-26 to verify last-two-slides duplication."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/MLA_Workshop_061X_Works_Cited.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_count = len(prs.slides)
            result = {'slide_count': slide_count}
            for idx, key in [(24, 'slide_25_text'), (25, 'slide_26_text')]:
                if 0 <= idx < slide_count:
                    texts = []
                    for shape in prs.slides[idx].shapes:
                        if shape.has_text_frame:
                            t = shape.text_frame.text.strip()
                            if t:
                                texts.append(t)
                    result[key] = ' | '.join(texts)
                else:
                    result[key] = ''
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_stretch_and_transition__676dd2418e7e3a84c9083aa5ad9f76e2_qw35sft2_861ec0f5(env, config: dict):
    import tempfile, os
    from pptx import Presentation
    from pptx.oxml.ns import qn
    pptx_path = config.get('path', '/home/user/Desktop/CPD_Background_Investigation_Process.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        slide_width_cm = prs.slide_width / 914400 * 2.54
        slide_height_cm = prs.slide_height / 914400 * 2.54
        best_shape = None
        best_width = 0.0
        for shape in slide.shapes:
            if shape.shape_type == 13 and 'Picture 2' in shape.name:
                w = shape.width / 914400 * 2.54
                if w > best_width:
                    best_width = w
                    best_shape = shape
        if best_shape is None:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    w = shape.width / 914400 * 2.54
                    if w > best_width:
                        best_width = w
                        best_shape = shape
        image_width_cm = best_shape.width / 914400 * 2.54 if best_shape else 0.0
        image_height_cm = best_shape.height / 914400 * 2.54 if best_shape else 0.0
        image_left_cm = best_shape.left / 914400 * 2.54 if best_shape else -1.0
        image_top_cm = best_shape.top / 914400 * 2.54 if best_shape else -1.0
        transition_type = 'none'
        transition_elem = slide._element.find(qn('p:transition'))
        if transition_elem is not None:
            for child in transition_elem:
                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if tag not in ('sndAc',):
                    transition_type = tag
                    break
            if transition_type == 'none' and len(list(transition_elem)) == 0:
                transition_type = 'cut'
        return {'image_width_cm': round(image_width_cm, 2), 'image_height_cm': round(image_height_cm, 2), 'image_left_cm': round(image_left_cm, 2), 'image_top_cm': round(image_top_cm, 2), 'slide_width_cm': round(slide_width_cm, 2), 'slide_height_cm': round(slide_height_cm, 2), 'transition_type': transition_type}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_info__881be08815c16b0f7ba88245e2d69e10_qw35sft2_207240ca(env, config: dict):
    """Get slide count and text frame info from a PPTX presentation on the VM."""
    import tempfile, os
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/slideshow.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': f'File not found: {path}'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        slides_with_text_frames = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slides_with_text_frames += 1
                    break
        return {'slide_count': slide_count, 'slides_with_text_frames': slides_with_text_frames}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_audio_and_transition__1b19ca869f4a35a02b8d110c7e17b7e9_qw35sft2_6f055a7e(env, config: dict):
    """Check if slide 1 has embedded audio and what transition type is set."""
    ppt_path = config.get('path', '/home/user/Desktop/Mady_and_Mia_Baseball.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(ppt_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_audio': False, 'transition_type': None}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_audio = False
        transition_type = None
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if rels_file in zf.namelist():
                with zf.open(rels_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    for rel in root.findall('r:Relationship', ns):
                        if 'audio' in rel.attrib.get('Type', ''):
                            has_audio = True
                            break
            slide_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            if slide_file in zf.namelist():
                with zf.open(slide_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    pns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    trans = root.find('.//{%s}transition' % pns)
                    if trans is not None:
                        children = list(trans)
                        if children:
                            child_tag = children[0].tag
                            transition_type = child_tag.split('}')[-1] if '}' in child_tag else child_tag
                        else:
                            transition_type = 'cut'
        return {'has_audio': has_audio, 'transition_type': transition_type, 'slide_index': slide_index}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_font_italic__5e9f703dc627e70a3b4d09452ec3d3ab_qw35sft2_b38d45fc(env, config: dict):
    """Get font name and italic state of the last slide title in 24_8.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/24_8.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[-1]
        title_shape = slide.shapes.title
        if title_shape is None or not title_shape.has_text_frame:
            return {'font_name': None, 'italic': None}
        font_names = []
        italic_values = []
        for para in title_shape.text_frame.paragraphs:
            para_italic = para.font.italic
            for run in para.runs:
                if run.text.strip():
                    font_names.append(run.font.name)
                    run_italic = run.font.italic
                    if run_italic is None:
                        run_italic = para_italic
                    italic_values.append(run_italic)
        if not font_names:
            return {'font_name': None, 'italic': None}
        non_none_names = [n for n in font_names if n is not None]
        font_name = non_none_names[0] if non_none_names else font_names[0]
        if all((v is True for v in italic_values)):
            italic = True
        elif any((v is False for v in italic_values)):
            italic = False
        else:
            italic = None
        return {'font_name': font_name, 'italic': italic}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slides35_colors__0fbb4b3515e1009c69bc08919cd2b491_qw35sft2_1cbad4a0(env, config: dict):
    """Get text run font colors from slides 3 and 5 of 1_2.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/1_2.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        result = {}
        for slide_idx, slide_key in [(2, 'slide3'), (4, 'slide5')]:
            slide = prs.slides[slide_idx]
            colors = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                if run.font.color and run.font.color.type is not None:
                                    colors.append(str(run.font.color.rgb).upper())
                                else:
                                    colors.append(None)
            result[slide_key] = colors
        return result
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide3_subpoint_level__3c50fb39e50d5c4bd155f8f4b7e911d1_qw35sft2_170ff70f(env, config: dict):
    """Read indentation level of 'first point of sub topics' on slide 3."""
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/Writing-Outlines.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[2]
        shape = slide.shapes[1]
        tf = shape.text_frame
        for para in tf.paragraphs:
            if 'first point of sub topics' in para.text:
                return {'level': para.level, 'text': para.text}
        return {'error': 'first point of sub topics paragraph not found'}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_file_slide_count__e4d3506bfaa42b2028f1a93213c30bc6_qw35sft2_f8e2a094(env, config: dict):
    """Get state of pre.pptx on Desktop: file existence and slide count."""
    import tempfile, os
    from pptx import Presentation
    desktop_path = '/home/user/Desktop/pre.pptx'
    file_bytes = env.controller.get_file(desktop_path)
    if not file_bytes:
        return {'file_exists': False, 'slide_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(list(prs.slides))
        return {'file_exists': True, 'slide_count': slide_count}
    finally:
        os.unlink(tmp_path)

def get_impress_bullet_underline__396e098239a8be5c53a96d982334cb7b_qw35sft2_75a01ec0(env, config: dict):
    """Get bullet status and underline status of content paragraph in 69_4.pptx."""
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/69_4.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        content_shape = slide.shapes[1]
        tf = content_shape.text_frame
        para = tf.paragraphs[0]
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        pPr = para._p.find(f'{{{ns}}}pPr')
        has_bullet = False
        if pPr is not None:
            buChar = pPr.find(f'{{{ns}}}buChar')
            has_bullet = buChar is not None
        has_underline = False
        if para.runs:
            has_underline = all((run.font.underline is True for run in para.runs))
        return {'has_bullet': has_bullet, 'has_underline': has_underline}
    finally:
        os.unlink(tmp_path)

def get_impress_options_combo__43da862b24de1bb5002a875bb93d754e_qw35sft2_7a41a212(env, config: dict):
    """Read LibreOffice registrymodifications.xcu to check presenter console disabled and autosave interval."""
    xcu_path = '/home/user/.config/libreoffice/4/user/registrymodifications.xcu'
    file_bytes = env.controller.get_file(xcu_path)
    if not file_bytes:
        return {'error': 'XCU file not found', 'presenter_disabled': False, 'autosave_minutes': None}
    try:
        content = file_bytes.decode('utf-8') if isinstance(file_bytes, bytes) else str(file_bytes)
        pattern_presenter = 'PresenterScreenEnabled[^<]*<value>(false|true)</value>'
        match_presenter = re.search(pattern_presenter, content, re.IGNORECASE)
        if match_presenter:
            presenter_disabled = match_presenter.group(1).lower() == 'false'
        else:
            presenter_disabled = False
        pattern_autosave = 'AutoSaveTimeIntervall[^<]*<value>(\\d+)</value>'
        match_autosave = re.search(pattern_autosave, content, re.IGNORECASE)
        autosave_minutes = int(match_autosave.group(1)) if match_autosave else None
        return {'presenter_disabled': presenter_disabled, 'autosave_minutes': autosave_minutes}
    except Exception as e:
        return {'error': str(e), 'presenter_disabled': False, 'autosave_minutes': None}

def get_pptx_slide_orientation__1830f729a2ae0637260a8bca58b0c8cd_qw35sft2_29845848(env, config: dict):
    """Get the slide dimensions to determine orientation of the PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Emu
        file_path = config.get('path', '/home/user/Desktop/note-taking-strategies.pptx')
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            width = prs.slide_width
            height = prs.slide_height
            is_portrait = height > width
            return {'width': int(width), 'height': int(height), 'orientation': 'portrait' if is_portrait else 'landscape'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_img_pos__edae955a62ba6a9d5bff7ce774e3ee10_qw35sft2_c281af6e(env, config: dict):
    """Get the top position (in EMU) of the picture shape on slide 2 of the pptx file."""
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/43_1.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slides = list(prs.slides)
        if len(slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide2 = slides[1]
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return {'image_top_emu': shape.top, 'image_left_emu': shape.left, 'slide_height_emu': prs.slide_height}
        return {'error': 'No picture found on slide 2'}
    finally:
        os.unlink(tmp_path)

def get_impress_title_color_underline__84ce98fac5a701a1ff43f1337a4f838c_qw35sft2_8e6a55dc(env, config: dict):
    """
    Get both font color and underline state of the topmost text shape (title)
    in a specified slide of a PPTX file (variation 4, slide 3).
    config keys: path (VM path), slide_index (0-based)
    Returns: {'color_hex': '000000', 'underline': True, 'title_text': ..., 'error': None}
    """
    from pptx import Presentation
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/4_1.pptx'))
    if not file_bytes:
        return {'error': 'File not found', 'color_hex': None, 'underline': None}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 2)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range', 'color_hex': None, 'underline': None}
        slide = prs.slides[slide_index]
        shapes_with_text = [(shape.top if shape.top is not None else float('inf'), shape) for shape in slide.shapes if hasattr(shape, 'text_frame') and shape.text.strip()]
        if not shapes_with_text:
            return {'error': 'No text shapes found on slide', 'color_hex': None, 'underline': None}
        shapes_with_text.sort(key=lambda x: x[0])
        title_shape = shapes_with_text[0][1]
        colors = []
        underlines = []
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                if font.color and font.color.type is not None:
                    try:
                        colors.append(str(font.color.rgb).upper())
                    except Exception:
                        colors.append(None)
                else:
                    colors.append(None)
                underlines.append(font.underline)
        color_val = colors[0] if colors else None
        all_underlined = all((u is True for u in underlines)) if underlines else False
        return {'color_hex': color_val, 'all_colors': colors, 'underline': all_underlined, 'all_underlines': underlines, 'title_text': title_shape.text.strip()[:80], 'error': None}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide2_title_bold__ba024eb57784584ce2b08e31db5871ac_qw35sft2_538afce5(env, config: dict):
    """Get slide 2 title text, alignment, and bold formatting from the PPTX file."""
    import tempfile, os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/22_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide2 = prs.slides[1]
        title_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                title_shape = shape
                break
        if title_shape is None:
            for shape in slide2.shapes:
                if shape.has_text_frame:
                    title_shape = shape
                    break
        if title_shape is None:
            return {'title_text': '', 'title_alignment': None, 'title_bold': False}
        tf = title_shape.text_frame
        title_text = tf.text.strip()
        alignment = None
        if tf.paragraphs:
            alignment = tf.paragraphs[0].alignment
        alignment_name = alignment.name if alignment is not None else None
        bold = False
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.bold:
                    bold = True
                    break
        return {'title_text': title_text, 'title_alignment': alignment_name, 'title_bold': bold}
    finally:
        os.unlink(tmp_path)

def get_pptx_image_position__2e5aa88aaf11fa4bda2a55acfcab9dc5_qw35sft2_9a585c0e(env, config: dict):
    """Get image position info from Slide 2 of 201_6.pptx."""
    file_path = config.get('path', '/home/user/Desktop/201_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide_width = prs.slide_width
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if shape.shape_type == 13:
                return {'image_left': shape.left, 'image_top': shape.top, 'image_width': shape.width, 'image_height': shape.height, 'slide_width': slide_width, 'slide_height': prs.slide_height}
        return {'error': 'No picture found on Slide 2'}
    finally:
        os.unlink(tmp_path)

def get_slide2_title_state__78dfac329c5cd8b1c5efec5b79600fdd_qw35sft2_c4a7bb63(env, config: dict):
    """Get bold state of the title on slide 2 of 164_3.pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'pptx module not available'}
    file_path = config.get('path', '/home/user/Desktop/164_3.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide = prs.slides[1]
        title_bold = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                continue
            if shape.placeholder_format.idx != 0:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        if run.font.bold is True:
                            title_bold = True
                        break
                try:
                    if para.font.bold is True:
                        title_bold = True
                except Exception:
                    pass
            break
        return {'title_bold': title_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_slide3_title_bold__7dcab435f4cba9d9e23699eed9a4d408_qw35sft2_11c65a9e(env, config: dict):
    """Get bold status of the title text in slide 3 of 21_0.pptx."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not installed'}
    file_path = config.get('path', '/home/user/Desktop/21_0.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[2]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if 'Add an Agenda Page' in text:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            for run in para.runs:
                                return {'title_text': text, 'title_bold': run.font.bold}
        return {'error': 'Title shape not found on slide 3', 'title_bold': None}
    finally:
        os.unlink(tmp_path)

def get_pptx_text_alignments__14ff9da6ca84bfee11448e7c7f8efafc_qw35sft2_5c719df5(env, config: dict):
    """Get text alignment for first textboxes on slides 3, 4, and 5 of 38_1.pptx."""
    import tempfile, os
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    file_path = config.get('path', '/home/user/Desktop/38_1.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        results = {}
        for key, slide_idx, keyword in [('slide3_align', 2, 'HELLO'), ('slide4_align', 3, 'WRITE YOUR'), ('slide5_align', 4, 'WRITE AN ORIGINAL')]:
            slide = prs.slides[slide_idx]
            found = None
            for shape in slide.shapes:
                para = None
                if hasattr(shape, 'text_frame') and keyword.lower() in shape.text_frame.text.lower():
                    para = shape.text_frame.paragraphs[0]
                elif shape.shape_type == 6:
                    for sub in shape.shapes:
                        if hasattr(sub, 'text_frame') and keyword.lower() in sub.text_frame.text.lower():
                            para = sub.text_frame.paragraphs[0]
                            break
                if para is not None:
                    align = para.alignment
                    if align == PP_ALIGN.RIGHT:
                        found = 'RIGHT'
                    elif align == PP_ALIGN.CENTER:
                        found = 'CENTER'
                    elif align == PP_ALIGN.LEFT:
                        found = 'LEFT'
                    else:
                        found = 'NONE'
                    break
            results[key] = found
        return results
    finally:
        os.unlink(tmp_path)

def get_impress_notes_bg_transition__5cb083f40f95f75316141f060ee8a1cc_qw35sft2_1fc3c009(env, config: dict):
    """Get slide notes, background RGB, and whether a transition is set from a PPTX."""
    import tempfile
    import os
    from pptx import Presentation
    from pptx.oxml.ns import qn
    file_path = config.get('path', '/home/user/Desktop/181_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        notes_text = ''
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
        bg_rgb = None
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            try:
                bg_rgb = str(fill.fore_color.rgb)
            except Exception:
                bg_rgb = None
        trans_elem = slide._element.find(qn('p:transition'))
        has_transition = trans_elem is not None
        transition_type = None
        if trans_elem is not None:
            for child in trans_elem:
                local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                transition_type = local
                break
        return {'notes_text': notes_text, 'bg_rgb': bg_rgb, 'has_transition': has_transition, 'transition_type': transition_type}
    finally:
        os.unlink(tmp_path)

def get_impress_slide14_font_bold__a28710b7c875e09d23efa44313d2c747_qw35sft2_4643d71f(env, config: dict):
    """Get font sizes and bold flag of the first two non-empty textboxes on slide 14."""
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/45_1.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[13]
        boxes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not tf.text.strip():
                continue
            size_pt = None
            bold = None
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        if run.font.size is not None:
                            size_pt = round(run.font.size / 12700)
                        bold = run.font.bold
                        break
                if size_pt is not None:
                    break
            boxes.append({'text': tf.text[:50], 'size_pt': size_pt, 'bold': bold})
            if len(boxes) == 2:
                break
        return {'textbox1_size_pt': boxes[0]['size_pt'] if len(boxes) > 0 else None, 'textbox1_bold': boxes[0]['bold'] if len(boxes) > 0 else None, 'textbox1_text': boxes[0]['text'] if len(boxes) > 0 else None, 'textbox2_size_pt': boxes[1]['size_pt'] if len(boxes) > 1 else None, 'textbox2_text': boxes[1]['text'] if len(boxes) > 1 else None}
    finally:
        os.unlink(tmp_path)

def get_impress_slide3_table_and_title__80d9e8d7efd024cbc51991913269b5fb_qw35sft2_53d7293b(env, config: dict):
    """Get table position and title text of slide 3 from the PPTX file."""
    import tempfile
    import os
    pptx_path = config.get('path', '/home/user/Desktop/55_10.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        if len(prs.slides) < 3:
            return {'error': f'Expected at least 3 slides, got {len(prs.slides)}'}
        slide3 = prs.slides[2]
        slide_height = prs.slide_height
        table_shape = None
        title_text = None
        for shape in slide3.shapes:
            if shape.shape_type == 19:
                table_shape = shape
            if shape.has_text_frame and hasattr(shape, 'placeholder_format') and (shape.placeholder_format is not None):
                if shape.placeholder_format.idx == 0:
                    title_text = shape.text_frame.text.strip()
        if table_shape is None:
            return {'error': 'No table found on slide 3', 'title': title_text}
        return {'table_top': table_shape.top, 'table_height': table_shape.height, 'slide_height': slide_height, 'title': title_text}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_title_full__9620e3371e8842ebb86c9c7ef6cb4f00_qw35sft2_c6d92324(env, config: dict):
    """Get the title text, font size, and font color of a specific slide in a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/189_4.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 1)
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame and shape.name == 'PlaceHolder 1':
                    text = shape.text_frame.text.strip()
                    font_size_pt = None
                    color_hex = None
                    r, g, b = (None, None, None)
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if font_size_pt is None and run.font.size is not None:
                                font_size_pt = round(run.font.size / 12700, 1)
                            if color_hex is None and run.font.color and (run.font.color.type is not None):
                                try:
                                    rgb = run.font.color.rgb
                                    color_hex = str(rgb).upper()
                                    r = int(color_hex[0:2], 16)
                                    g = int(color_hex[2:4], 16)
                                    b = int(color_hex[4:6], 16)
                                except Exception:
                                    pass
                        if font_size_pt is not None and color_hex is not None:
                            break
                    return {'title_text': text, 'font_size_pt': font_size_pt, 'color_hex': color_hex, 'r': r, 'g': g, 'b': b}
            return {'error': 'Title placeholder not found', 'title_text': '', 'font_size_pt': None, 'color_hex': None}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_title_font_transition__7950498f55dec024fcd49c95968dd860_qw35sft2_a7a84401(env, config: dict):
    """Get title text, font name, and slide transition type from slide 1 of 9_1.pptx."""
    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/9_1.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            slide = prs.slides[slide_index]
            title_text = None
            font_name = None
            transition_type = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt and title_text is None:
                        title_text = txt
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    font_name = run.font.name
                                break
                            if font_name:
                                break
            slide_elem = slide._element
            transition_elem = slide_elem.find(qn('p:transition'))
            if transition_elem is not None:
                for trans_name in ['fade', 'blinds', 'checker', 'circle', 'dissolve', 'push', 'wipe', 'zoom', 'wheel', 'strips', 'fly', 'split', 'wedge', 'newsflash', 'random']:
                    child = transition_elem.find(qn(f'p:{trans_name}'))
                    if child is not None:
                        transition_type = trans_name
                        break
                if transition_type is None:
                    transition_type = 'unknown'
            return {'title_text': title_text, 'font_name': font_name, 'transition_type': transition_type}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_text_color__f2292754a9e180ffad3521e45a75c905_qw35sft2_d84f0cfa(env, config: dict):
    """Get the font color of a specific text on a slide in the PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/saa-format-guide.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            target_text = config.get('target_text', '')
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.replace('\n', ' ').strip()
                if target_text and target_text.upper() not in text.upper():
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.type is not None:
                            try:
                                rgb = run.font.color.rgb
                                hex_str = str(rgb)
                                return {'text': text, 'color_hex': hex_str.upper(), 'r': int(hex_str[0:2], 16), 'g': int(hex_str[2:4], 16), 'b': int(hex_str[4:6], 16)}
                            except Exception:
                                pass
            return {'error': f'Text "{target_text}" not found or has no explicit color set'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_slide2_bg_title_bold__3ee304ba526b814e3d4d6c2be946eae0_qw35sft2_71d9d6fc(env, config: dict):
    """Get slide 2 background color and whether the title text is bold from 71_6.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/71_6.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[1]
        bg_color = None
        try:
            fill = slide.background.fill
            if fill.type is not None and int(fill.type) == 1:
                bg_color = str(fill.fore_color.rgb).upper()
        except Exception:
            pass
        title_bold = None
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        if text_shapes:
            title_shape = text_shapes[0]
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    title_bold = run.font.bold
                    break
                if title_bold is not None:
                    break
        return {'bg_color': bg_color, 'title_bold': title_bold}
    finally:
        os.unlink(tmp_path)

def get_pptx_single_text_color__c8d6e08f7da191ce71fc333f5a053d63_qw35sft2_f61679a0(env, config: dict):
    """Get the font color of a specific textbox in slide 1 of a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/45_2.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            target_text = config.get('target_text', 'NEW PR')
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.replace('\n', ' ').strip()
                    if target_text.upper() in text.upper():
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.type is not None:
                                    try:
                                        rgb = run.font.color.rgb
                                        hex_str = str(rgb)
                                        r = int(hex_str[0:2], 16)
                                        g = int(hex_str[2:4], 16)
                                        b = int(hex_str[4:6], 16)
                                        return {'text': text, 'color_hex': hex_str, 'r': r, 'g': g, 'b': b}
                                    except Exception:
                                        pass
            return {'error': f'Text "{target_text}" not found or has no explicit color'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_all_text_color__2c3419decc4324f20da8cbd8ae0e6c2d_qw35sft2_90378413(env, config: dict):
    """Get font color status for all text shapes (title, body, table) on slide 1."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/154_3.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide = prs.slides[0]
            expected_hex = config.get('expected_color', 'C00000').upper().lstrip('#')
            tolerance = config.get('tolerance', 15)
            try:
                exp_r = int(expected_hex[0:2], 16)
                exp_g = int(expected_hex[2:4], 16)
                exp_b = int(expected_hex[4:6], 16)
            except Exception:
                return {'error': 'Invalid expected_color in config'}
            title_ok = False
            if slide.shapes[0].has_text_frame:
                title_runs = [r for p in slide.shapes[0].text_frame.paragraphs for r in p.runs]
                if title_runs:
                    title_ok = True
                    for run in title_runs:
                        try:
                            if run.font.color is None or run.font.color.type is None:
                                title_ok = False
                                break
                            h = str(run.font.color.rgb)
                            if not (abs(int(h[0:2], 16) - exp_r) <= tolerance and abs(int(h[2:4], 16) - exp_g) <= tolerance and (abs(int(h[4:6], 16) - exp_b) <= tolerance)):
                                title_ok = False
                                break
                        except Exception:
                            title_ok = False
                            break
            body_ok = False
            if slide.shapes[1].has_text_frame:
                body_runs = [r for p in slide.shapes[1].text_frame.paragraphs for r in p.runs]
                if body_runs:
                    body_ok = True
                    for run in body_runs:
                        try:
                            if run.font.color is None or run.font.color.type is None:
                                body_ok = False
                                break
                            h = str(run.font.color.rgb)
                            if not (abs(int(h[0:2], 16) - exp_r) <= tolerance and abs(int(h[2:4], 16) - exp_g) <= tolerance and (abs(int(h[4:6], 16) - exp_b) <= tolerance)):
                                body_ok = False
                                break
                        except Exception:
                            body_ok = False
                            break
            table_ok = False
            table_shape = slide.shapes[2]
            if table_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_runs = [run for row in table_shape.table.rows for cell in row.cells for para in cell.text_frame.paragraphs for run in para.runs]
                if table_runs:
                    table_ok = True
                    for run in table_runs:
                        try:
                            if run.font.color is None or run.font.color.type is None:
                                table_ok = False
                                break
                            h = str(run.font.color.rgb)
                            if not (abs(int(h[0:2], 16) - exp_r) <= tolerance and abs(int(h[2:4], 16) - exp_g) <= tolerance and (abs(int(h[4:6], 16) - exp_b) <= tolerance)):
                                table_ok = False
                                break
                        except Exception:
                            table_ok = False
                            break
            return {'title_color_ok': title_ok, 'body_color_ok': body_ok, 'table_color_ok': table_ok}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_content_bold__87b69c70091c302ea2d6ddb5f5d9c002_qw35sft2_eb05701d(env, config: dict):
    """Get Slide 2 content text and whether any run in that text is bold."""
    import tempfile
    import os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/109_4.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide2_content = ''
        slide2_content_bold = False
        if len(prs.slides) >= 2:
            slide2 = prs.slides[1]
            for shape in slide2.shapes:
                if shape.has_text_frame and shape.name == 'PlaceHolder 2':
                    slide2_content = shape.text_frame.text.strip()
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.bold:
                                slide2_content_bold = True
                    break
        return {'slide2_content': slide2_content, 'slide2_content_bold': slide2_content_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_slide_pane__0dbd76d4f35e1d6b60aa20d8163a50c2_qw35sft2_d838fcea(env, config: dict):
    """Get LibreOffice Impress UI state to check if the left slide pane is visible."""
    tree = env.controller.get_accessibility_tree()
    return {'tree': tree if isinstance(tree, str) else ''}

def get_pptx_text_font_size__8cd3b2c6f2a2d80e7aa650d3b0962695_qw35sft2_7904b228(env, config: dict):
    """Get font sizes of all runs in a specific shape of the first slide."""
    import tempfile, os
    from pptx import Presentation
    path = config.get('path', '/home/user/Desktop/16_2.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        shape_index = config.get('shape_index', 0)
        shape = slide.shapes[shape_index]
        if not shape.has_text_frame:
            return {'error': 'Shape has no text frame'}
        sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
        return {'font_sizes': sizes, 'shape_index': shape_index}
    finally:
        os.unlink(tmp_path)

def get_impress_pptx_props__e526a93f073489d810264dc88333d90d_qw35sft2_38da5688(env, config: dict):
    """Get slide 3 Group 6 height and slide 6 textbox font sizes from 42_2.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/42_2.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        result = {}
        slide3 = prs.slides[2]
        result['slide3_group6_height_cm'] = None
        for shape in slide3.shapes:
            if shape.name == 'Group 6':
                result['slide3_group6_height_cm'] = round(shape.height / 360000, 4)
                break
        slide6 = prs.slides[5]
        slide6_font_sizes = {}
        for shape in slide6.shapes:
            if shape.shape_type == 1 and hasattr(shape, 'text_frame') and shape.text.strip():
                pts = []
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            pts.append(int(round(run.font.size.pt)))
                if pts:
                    slide6_font_sizes[shape.name] = pts
        result['slide6_font_sizes'] = slide6_font_sizes
        return result
    finally:
        os.unlink(tmp_path)

def get_pptx_transition__24d25146d014952eed8e287afb1ec612_qw35sft2_f90ee5a0(env, config: dict):
    """Check whether slide 0 of a PPTX file has any transition applied."""
    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/AM_Last_Page_Template.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 0)
            if slide_index >= len(prs.slides):
                return {'error': 'Slide index out of range'}
            slide = prs.slides[slide_index]
            transition_el = slide._element.find(qn('p:transition'))
            if transition_el is None:
                return {'has_transition': False, 'transition_xml': None}
            return {'has_transition': True, 'transition_xml': transition_el.xml if hasattr(transition_el, 'xml') else str(transition_el.tag)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_transitions__b146bd229b2aa1513d722fe51dc492dc_qw35sft2_9710b95b(env, config: dict):
    """Read slide transitions for multiple slides from a PPTX file on the VM."""
    file_path = config.get('path', '/home/user/Desktop/Ch4 Video Effect.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        namespaces = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        result = {}
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            slide_count = sum((1 for name in zf.namelist() if name.startswith('ppt/slides/slide') and name.endswith('.xml')))
            result['slide_count'] = slide_count
            for idx in range(slide_count):
                slide_name = 'ppt/slides/slide{}.xml'.format(idx + 1)
                try:
                    with zf.open(slide_name) as sf:
                        tree = ET.parse(sf)
                        root = tree.getroot()
                        transition = root.find('.//p:transition', namespaces)
                        if transition is not None:
                            children = list(transition)
                            if children:
                                tag = children[0].tag
                                local = tag.split('}')[-1] if '}' in tag else tag
                                result['slide_{}'.format(idx)] = local
                            else:
                                result['slide_{}'.format(idx)] = 'set_no_type'
                        else:
                            result['slide_{}'.format(idx)] = None
                except Exception:
                    result['slide_{}'.format(idx)] = None
        return result
    finally:
        os.unlink(tmp_path)

def get_pptx_table_row0__a30c5776ecd0bc922b282656ddbb8d0e_qw35sft2_1f43e00d(env, config: dict):
    """Get the first row of a table on a specific slide from a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/33_1.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 3)
            slide = prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.shape_type == 19:
                    row0 = shape.table.rows[0]
                    table_row0 = [cell.text.strip() for cell in row0.cells]
                    return {'table_row0': table_row0}
            return {'error': 'No table found on specified slide'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_title_text__a33045396be2560e0d47268bd7a68b1f_qw35sft2_80ff8e81(env, config: dict):
    """Get the title placeholder text from a specific slide in a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/214_9.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_idx = config.get('slide_index', 0)
            if slide_idx >= len(prs.slides):
                return {'error': f'Slide index {slide_idx} out of range'}
            slide = prs.slides[slide_idx]
            title_text = ''
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    pf = shape.placeholder_format
                    if pf is not None and pf.type.real == 1:
                        title_text = shape.text_frame.text.strip()
                        break
                except Exception:
                    continue
            return {'title_text': title_text}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide_bg_color__3e8a909d4fdee8dfddabe4251505a445_qw35sft2_a004fe09(env, config: dict):
    """Get background color of a specific slide in a PPTX file.

    Returns dict with fill_type, rgb (hex string, 6 chars), and fallback
    master_rgb when the slide itself does not have an explicit solid fill.
    """
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/lec17-gui-events.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'rgb': None, 'master_rgb': None}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from pptx import Presentation
        prs = Presentation(tmp_path)
        slide_index = config.get('slide_index', 0)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range (has {len(prs.slides)} slides)', 'rgb': None, 'master_rgb': None}
        slide = prs.slides[slide_index]
        bg = slide.background
        fill = bg.fill
        result = {'fill_type': fill.type.name if fill.type is not None else 'NONE', 'rgb': None, 'master_rgb': None, 'slide_index': slide_index, 'num_slides': len(prs.slides)}
        if fill.type is not None and fill.type.name == 'SOLID':
            try:
                rgb = fill.fore_color.rgb
                result['rgb'] = str(rgb).upper()
            except Exception as e:
                result['rgb_error'] = str(e)
        if result['rgb'] is None and len(prs.slide_masters) > 0:
            master_fill = prs.slide_masters[0].background.fill
            result['master_fill_type'] = master_fill.type.name if master_fill.type is not None else 'NONE'
            if master_fill.type is not None and master_fill.type.name == 'SOLID':
                try:
                    rgb = master_fill.fore_color.rgb
                    result['master_rgb'] = str(rgb).upper()
                except Exception as e:
                    result['master_rgb_error'] = str(e)
        return result
    except Exception as e:
        return {'error': str(e), 'rgb': None, 'master_rgb': None}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_impress_slide_info__b298657a7c2d705873344441c4e6373c_qw35sft2_3775b7a2(env, config: dict):
    """Get slide count and text frame info from a PPTX presentation on the VM."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/three_slides.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': f'File not found: {path}'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide_count = len(prs.slides)
        slides_with_text_frames = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slides_with_text_frames += 1
                    break
        return {'slide_count': slide_count, 'slides_with_text_frames': slides_with_text_frames}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_pptx_slide_image_and_title__68b16c8fa261fe6f7f7bf99cd53e207c_qw35sft2_7de26c83(env, config: dict):
    """Get image info and title text from a specified slide in a PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/31_2.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_index = config.get('slide_index', 1)
            if slide_index >= len(prs.slides):
                return {'error': f'Slide index {slide_index} out of range, total slides: {len(prs.slides)}'}
            slide = prs.slides[slide_index]
            images = []
            title_text = None
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append({'name': shape.name, 'width_cm': round(shape.width / 360000, 4), 'height_cm': round(shape.height / 360000, 4)})
                elif hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    try:
                        if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                            title_text = shape.text_frame.text.strip()
                    except Exception:
                        pass
            return {'image_count': len(images), 'images': images, 'title_text': title_text}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_impress_slide2_title_props__8ff66c4d1cd7e4775354b5fdf25e12eb_qw35sft2_fc2480c5(env, config: dict):
    """Get title placeholder position and font size on slide 2 of 134_2.pptx."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not installed'}
    file_path = config.get('path', '/home/user/Desktop/134_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if len(prs.slides) < 2:
            return {'error': 'Slide 2 not found'}
        slide = prs.slides[1]
        title_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                title_shape = shape
                break
        if title_shape is None:
            return {'error': 'Title placeholder not found on slide 2'}
        top_cm = title_shape.top / 360000.0
        font_size_pt = None
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    font_size_pt = run.font.size.pt
                    break
            if font_size_pt is not None:
                break
        if font_size_pt is None:
            for para in title_shape.text_frame.paragraphs:
                if para.font.size is not None:
                    font_size_pt = para.font.size.pt
                    break
        return {'top_cm': top_cm, 'font_size_pt': font_size_pt}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_pptx_notes__e6eecf5158b3e658cb4aa48b9b009898_qw35sft2_40422256(env, config: dict):
    """Get the notes text of slide 0 from the PPTX file."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    path = config.get('path', '/home/user/Desktop/186_3.pptx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        if not slide.has_notes_slide:
            return {'notes': ''}
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        return {'notes': notes_text}
    finally:
        os.unlink(tmp_path)

def get_pptx_two_slide_tables__422453ea9133a41f5d3b73c5d61c25f4_qw35sft2_68f346c8(env, config: dict):
    """Get table states from two slides (Product Overview and Features) in the PPTX file."""
    import tempfile, os
    try:
        from pptx import Presentation
    except ImportError:
        return {'error': 'python-pptx not available'}
    file_path = config.get('path', '/home/user/Desktop/41_3.pptx')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            total_slides = len(prs.slides)
            slide2_tables = []
            if 1 < total_slides:
                for shape in prs.slides[1].shapes:
                    if shape.has_table:
                        tbl = shape.table
                        slide2_tables.append({'rows': len(tbl.rows), 'cols': len(tbl.columns)})
            slide3_tables = []
            if 2 < total_slides:
                for shape in prs.slides[2].shapes:
                    if shape.has_table:
                        tbl = shape.table
                        slide3_tables.append({'rows': len(tbl.rows), 'cols': len(tbl.columns)})
            return {'total_slides': total_slides, 'slide2_tables': slide2_tables, 'slide3_tables': slide3_tables}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_slide5_text_colors__7491f6c65f733a38f893bdf8cad068b0_qw35sft2_48da877c(env, config: dict):
    """Get all text run font colors from slide 5 of 1_2.pptx."""
    import tempfile, os
    from pptx import Presentation
    file_bytes = env.controller.get_file('/home/user/Desktop/1_2.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[4]
        colors = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            if run.font.color and run.font.color.type is not None:
                                colors.append(str(run.font.color.rgb).upper())
                            else:
                                colors.append(None)
        return {'colors': colors, 'count': len(colors)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_impress_has_audio__5688a332db768af3543ab15e11555c76_qw35sft2_e375e116(env, config: dict):
    """Check if a specific slide in a PPTX file has an embedded audio file."""
    ppt_path = config.get('path', '/home/user/Desktop/Mady_and_Mia_Baseball.pptx')
    slide_index = config.get('slide_index', 0)
    file_bytes = env.controller.get_file(ppt_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_audio': False}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_audio = False
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if rels_file in zf.namelist():
                with zf.open(rels_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    for rel in root.findall('r:Relationship', ns):
                        if 'audio' in rel.attrib.get('Type', ''):
                            has_audio = True
                            break
        return {'has_audio': has_audio, 'slide_index': slide_index}
    finally:
        os.unlink(tmp_path)

def get_pptx_multi_slide_bg_colors__01d7b3b210af327bea300a303c9bddd6_qw35sft2_eb699e59(env, config: dict):
    """Get background solid fill colors for multiple slides in the PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/13_0.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide_indices = config.get('slide_indices', [2, 3])
            colors = {}
            for idx in slide_indices:
                if idx >= len(prs.slides):
                    colors[str(idx)] = {'error': f'Slide index {idx} out of range'}
                    continue
                slide = prs.slides[idx]
                bg = slide.background
                fill = bg.fill
                if fill.type is None:
                    colors[str(idx)] = {'color_hex': None, 'r': None, 'g': None, 'b': None}
                    continue
                try:
                    rgb = fill.fore_color.rgb
                    hex_str = str(rgb)
                    r = int(hex_str[0:2], 16)
                    g = int(hex_str[2:4], 16)
                    b = int(hex_str[4:6], 16)
                    colors[str(idx)] = {'color_hex': hex_str, 'r': r, 'g': g, 'b': b}
                except Exception as e:
                    colors[str(idx)] = {'error': str(e), 'color_hex': None}
            return {'slide_colors': colors}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pptx_stretch_and_text__04e2292403d24fd88b4576e3fc170228_qw35sft2_39c2c38f(env, config: dict):
    import tempfile, os
    from pptx import Presentation
    pptx_path = config.get('path', '/home/user/Desktop/CPD_Background_Investigation_Process.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        slide_width_cm = prs.slide_width / 914400 * 2.54
        slide_height_cm = prs.slide_height / 914400 * 2.54
        best_shape = None
        best_width = 0.0
        for shape in slide.shapes:
            if shape.shape_type == 13 and 'Picture 2' in shape.name:
                w = shape.width / 914400 * 2.54
                if w > best_width:
                    best_width = w
                    best_shape = shape
        image_width_cm = best_shape.width / 914400 * 2.54 if best_shape else 0.0
        image_height_cm = best_shape.height / 914400 * 2.54 if best_shape else 0.0
        image_left_cm = best_shape.left / 914400 * 2.54 if best_shape else -1.0
        image_top_cm = best_shape.top / 914400 * 2.54 if best_shape else -1.0
        textbox9_text = None
        for shape in slide.shapes:
            if shape.name == 'TextBox 9' and hasattr(shape, 'text'):
                textbox9_text = shape.text.strip()
                break
        bottom_text = None
        if textbox9_text is None:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text and ('Investigation' in shape.text):
                    bottom_text = shape.text.strip()
                    break
        return {'image_width_cm': round(image_width_cm, 2), 'image_height_cm': round(image_height_cm, 2), 'image_left_cm': round(image_left_cm, 2), 'image_top_cm': round(image_top_cm, 2), 'slide_width_cm': round(slide_width_cm, 2), 'slide_height_cm': round(slide_height_cm, 2), 'textbox9_text': textbox9_text, 'bottom_text_fallback': bottom_text}
    finally:
        os.unlink(tmp_path)

def get_lecture_slides_files__ac6b45f9bcb788e83171ced5593fc0f6_qw35sft2_8881d264(env, config: dict):
    """List files in /home/user/lecture_slides/ and check for target PDF."""
    result = env.controller.run_bash_script('ls /home/user/lecture_slides/ 2>/dev/null || echo "__EMPTY__"', timeout=15)
    if not result or result.get('returncode', -1) != 0:
        return {'error': 'Could not list lecture_slides directory', 'files': []}
    stdout = result.get('stdout', '').strip()
    if stdout == '__EMPTY__' or not stdout:
        return {'files': []}
    files = [f.strip() for f in stdout.split('\n') if f.strip()]
    return {'files': files}

def get_impress_video_docs_vlc__bae283564a19c9e043588bc4818877d3_qw35sft2_2a121f76(env, config: dict):
    """Check if a video file exists in ~/Documents/ and if VLC is running."""
    video_result = env.controller.run_bash_script('find /home/user/Documents -maxdepth 3 \\( -name "*.webm" -o -name "*.ogv" -o -name "*.mp4" \\) 2>/dev/null | head -5', timeout=30)
    video_output = ''
    if video_result and isinstance(video_result, dict):
        video_output = video_result.get('output', '').strip()
    video_files = [f.strip() for f in video_output.split('\n') if f.strip()] if video_output else []
    vlc_result = env.controller.run_bash_script('pgrep -x vlc 2>/dev/null || pgrep -x VLC 2>/dev/null || echo ""', timeout=15)
    vlc_output = ''
    if vlc_result and isinstance(vlc_result, dict):
        vlc_output = vlc_result.get('output', '').strip()
    vlc_running = bool(vlc_output)
    return {'video_in_documents': len(video_files) > 0, 'video_files': video_files, 'vlc_running': vlc_running}

def get_impress_compose4__c60351f10d67b0cf97ff8e001ca27d87_qw35sft2_5f92eccc(env, config: dict):
    """
    Check if background.png exists on Desktop AND whether the saved PPTX has
    a brightness (a:lum) attribute on the slide 2 background image element,
    indicating the LibreOffice brightness setting was saved into the file.
    """
    import tempfile
    import os
    result = {'file_exists': False, 'pptx_has_brightness': False, 'brightness_value': None, 'error': None}
    try:
        bg_bytes = env.controller.get_file('/home/user/Desktop/background.png')
        result['file_exists'] = bg_bytes is not None and len(bg_bytes) > 0
    except Exception as e:
        result['error'] = f'file_check error: {e}'
    try:
        pptx_bytes = env.controller.get_file('/home/user/Desktop/PPT-Template_widescreen.pptx')
        if pptx_bytes:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp.write(pptx_bytes)
                tmp_path = tmp.name
            try:
                from pptx import Presentation
                from lxml import etree
                prs = Presentation(tmp_path)
                slide2 = prs.slides[1]
                ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                for shape in slide2.shapes:
                    if shape.shape_type == 13:
                        elem = shape._element
                        lum = elem.find(f'.//{{{ns_a}}}lum')
                        if lum is not None:
                            result['pptx_has_brightness'] = True
                            bright_attr = lum.get('bright', None)
                            if bright_attr is not None:
                                result['brightness_value'] = int(bright_attr) / 1000.0
                            break
            finally:
                os.unlink(tmp_path)
    except Exception as e:
        result['error'] = (result.get('error') or '') + f' | pptx_check error: {e}'
    return result

def get_pptx_slide2_bg_and_count__24ae1553823b67948a98514a9ed74ad3_qw35sft2_e8ecc775(env, config: dict):
    """Check slide 2 background and verify total slide count is preserved."""
    import tempfile, os, zipfile, re
    pptx_path = config.get('path', '/home/user/Desktop/Robotic_Workshop_Infographics.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_background_image': False, 'slide_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as z:
            names = z.namelist()
            slide_files = [n for n in names if re.match('ppt/slides/slide\\d+\\.xml$', n)]
            slide_count = len(slide_files)
            slide2_has_bg = False
            if 'ppt/slides/slide2.xml' in names:
                slide2_xml = z.read('ppt/slides/slide2.xml').decode('utf-8')
                slide2_has_bg = '<p:bg>' in slide2_xml and '<a:blipFill' in slide2_xml
            slide1_has_new_bg = False
            if 'ppt/slides/slide1.xml' in names:
                slide1_xml = z.read('ppt/slides/slide1.xml').decode('utf-8')
                slide1_has_new_bg = '<p:bg>' in slide1_xml and '<a:blipFill' in slide1_xml
        return {'slide_count': slide_count, 'has_background_image': slide2_has_bg, 'slide1_unchanged': not slide1_has_new_bg}
    except Exception as e:
        return {'error': str(e), 'has_background_image': False, 'slide_count': 0}
    finally:
        os.unlink(tmp_path)

def get_lecture_slides_files__6f1f6989222847d5b9f522129093f325_qw35sft2_5e98a87d(env, config: dict):
    """List files in /home/user/lecture_slides/ to verify PDF download."""
    result = env.controller.run_bash_script('ls /home/user/lecture_slides/ 2>/dev/null || echo "__EMPTY__"', timeout=15)
    if not result or result.get('returncode', -1) != 0:
        return {'error': 'Could not list lecture_slides directory', 'files': []}
    stdout = result.get('stdout', '').strip()
    if stdout == '__EMPTY__' or not stdout:
        return {'files': []}
    files = [f.strip() for f in stdout.split('\n') if f.strip()]
    return {'files': files}

def get_pptx_notes_and_docx__1f2059f6b413a703d90fb8946ec0e698_qw35sft2_7af43bbc(env, config: dict):
    """
    Check both:
    1. The PPTX file for slide 7's note text
    2. The docx file for extracted notes including slide 7
    """
    pptx_path = config.get('pptx_path', '/home/user/Desktop/Dickinson_Slides.pptx')
    docx_path = config.get('docx_path', '/home/user/Desktop/notes.docx')
    result = {'pptx_slide7_note': None, 'docx_lines': [], 'errors': []}
    pptx_bytes = env.controller.get_file(pptx_path)
    if not pptx_bytes:
        result['errors'].append('PPTX not found')
    else:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(pptx_bytes)
            tmp_path = tmp.name
        try:
            from pptx import Presentation
            prs = Presentation(tmp_path)
            if len(prs.slides) >= 7:
                slide7 = prs.slides[6]
                if slide7.has_notes_slide:
                    note_text = slide7.notes_slide.notes_text_frame.text.strip()
                    result['pptx_slide7_note'] = note_text
                else:
                    result['pptx_slide7_note'] = ''
        except Exception as e:
            result['errors'].append(f'PPTX error: {e}')
        finally:
            os.unlink(tmp_path)
    docx_bytes = env.controller.get_file(docx_path)
    if not docx_bytes:
        result['errors'].append('docx not found')
    else:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(docx_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            lines = [para.text for para in doc.paragraphs if para.text.strip()]
            result['docx_lines'] = lines
        except Exception as e:
            result['errors'].append(f'docx error: {e}')
        finally:
            os.unlink(tmp_path)
    return result

def get_pptx_slide2_bg__d5c6210d63b8dbe359320c4bfae3310a_qw35sft2_ed30e2b7(env, config: dict):
    """Check whether slide 2 of the presentation has a picture background."""
    import tempfile, os, zipfile, re
    pptx_path = config.get('path', '/home/user/Desktop/Robotic_Workshop_Infographics.pptx')
    file_bytes = env.controller.get_file(pptx_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_background_image': False}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as z:
            if 'ppt/slides/slide2.xml' not in z.namelist():
                return {'error': 'slide2.xml not found in PPTX', 'has_background_image': False}
            slide2_xml = z.read('ppt/slides/slide2.xml').decode('utf-8')
        has_bg_elem = '<p:bg>' in slide2_xml
        has_blip_fill = '<a:blipFill' in slide2_xml and has_bg_elem
        has_solid_fill = '<a:solidFill' in slide2_xml and has_bg_elem
        has_any_fill = has_blip_fill or has_solid_fill
        return {'has_background_image': has_blip_fill, 'has_any_bg_fill': has_any_fill, 'has_bg_element': has_bg_elem, 'slide2_xml_length': len(slide2_xml)}
    except Exception as e:
        return {'error': str(e), 'has_background_image': False}
    finally:
        os.unlink(tmp_path)

def get_impress_compose0__9880fcd0ac1db0e1d57c564194f1780a_qw35sft2_622ef45d(env, config: dict):
    """
    Check if background.png exists on the Desktop, verify image brightness (~50%) in PPTX
    XML (<a:lum bright>), and get slide 2 title from saved PPTX.
    """
    import tempfile
    import os
    result = {'file_exists': False, 'brightness_ok': False, 'slide2_title': None, 'error': None}
    try:
        bg_bytes = env.controller.get_file('/home/user/Desktop/background.png')
        result['file_exists'] = bg_bytes is not None and len(bg_bytes) > 0
    except Exception as e:
        result['error'] = f'file_check error: {e}'
    try:
        pptx_bytes = env.controller.get_file('/home/user/Desktop/PPT-Template_widescreen.pptx')
        if pptx_bytes:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp.write(pptx_bytes)
                tmp_path = tmp.name
            try:
                from pptx import Presentation
                prs = Presentation(tmp_path)
                slide2 = prs.slides[1]
                slide2_title = None
                title_shape = slide2.shapes.title
                if title_shape is not None and title_shape.has_text_frame:
                    text = title_shape.text_frame.text.strip()
                    if text:
                        slide2_title = text
                if slide2_title is None:
                    for shape in slide2.shapes:
                        if shape.has_text_frame:
                            ph = getattr(shape, 'placeholder_format', None)
                            if ph is not None and ph.idx == 0:
                                slide2_title = shape.text_frame.text.strip()
                                break
                if slide2_title is None:
                    for shape in slide2.shapes:
                        if shape.has_text_frame:
                            text = shape.text_frame.text.strip()
                            if text.upper() in ('SLIDE TITLE', 'ENHANCED SLIDE'):
                                slide2_title = text
                                break
                result['slide2_title'] = slide2_title
                ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                for shape in slide2.shapes:
                    if shape.shape_type in (13, 14):
                        blip = shape.element.find('.//{%s}blip' % ns)
                        if blip is not None:
                            lum_el = blip.find('{%s}lum' % ns)
                            if lum_el is not None:
                                bright_str = lum_el.get('bright', '0')
                                try:
                                    bright_pct = int(bright_str) / 1000.0
                                    result['brightness_ok'] = abs(bright_pct - 50.0) <= 5.0
                                except (ValueError, TypeError):
                                    pass
                        break
            finally:
                os.unlink(tmp_path)
    except Exception as e:
        result['error'] = (result.get('error') or '') + f' | pptx_check error: {e}'
    return result
