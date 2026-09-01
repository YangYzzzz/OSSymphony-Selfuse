"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_file_dual_location__454f85dd5a1086d0f2418681bdac8477', 'get_manifest_json__ab85f9bbe05040b98eb1a964cbee7866', 'get_odt_files_exist__016af9433419594b3dc790848ad51e3b', 'get_text_file_content__b150ff4c0a473c93c912a80df923233a', 'get_folder_file_list__448eb54c6f660c787448d046dde3cea3', 'get_file_move_check__e60640f858dbfcce19e818f08302a608', 'get_pdf_page_info__a8da1aa98de37fa62b8d321ad0589eaf', 'get_main_py_content__bef206759cc52e3fdba71a4f0e77e8b3', 'get_txt_file_lines__d445ea0218043b2b599f555fb0b0c923', 'get_text_file_content__ca2a0f6f5509af90b7b3cf524daf3368', 'get_file_content__4c4a9bd632c3f7487e991bbd8a3f942a', 'get_verified_folder_files__445a1e82d5445803b67308e02fdc3544', 'get_docx_to_pdf_state__8f355d6af58b179cf68cc37725d92b0c', 'get_pdf_page_count__ddacf8dd85e44a11a4923ff07ce0abff', 'get_file_content__29f89216bfdb55e864e9190350cbf19a', 'get_csv_file_check__02188a8f533477bfdcad341f72281666', 'get_user_homedir__df678257c48d0ebe8f2e03968592b783', 'get_text_file_content__204b2db58c306983ce47277dbc39788e', 'get_file_organize_state__a65eb6fd760478532b3e118be8d079f9', 'get_file_copy_and_clipboard__267da608b7e8dd72021770568c2da84f', 'get_bib_file_content__dd385e64b1e5f06bf4b4c074020cb6a0', 'get_song_list_file__09209b5893b5a25bd669841ecee9c784', 'get_docx_text_content__b98b6fc5ce8ebe003cf7eaeacc71d919', 'get_file_rename_check__9b65526186f9fb77c93ff9d7de5c5fcb', 'get_def_file__9f211d2fef2982f0959041bc806f53c4', 'get_dir_existence__93ffa0d39e3e07dab53be514a31da362', 'get_vm_file_content__d175f88bb1991b0561f10b03c2a108ce', 'get_vm_file_content__6c3110d7dd0b28f9da7376f5ce78a46b', 'get_music_folder_state__eed664169e1770df941a91e7673f0e98', 'get_targz_contents__d544b10f933e46c3b8f208241fda6c7e', 'get_pdf_list_file__f72fa74134cf5a58d1281947ac1e62bd', 'get_file_content__79ff096ee36ebe0503f424e88a332ca2', 'get_docx_text_content__5ac00db9c05416b4087c6a9ef6d0dee7', 'get_docx_content__2aeb38ddfad7ab2f7bd9e6f63f6cc1bf', 'get_ods_to_pdf_check__40ad56f2b3098a6972e48b06b58c95e0', 'get_folder_file_list__69659eabd73459db19e2b0d1637cc759', 'get_novel_info_file__6a9459e3699067e60bf00ce031fcbd18', 'get_bib_file_content__db9975908e0c2c40a4467f2add6ed446', 'get_vm_file', 'get_file_content__05446d05a2730db5bc45314af39c65c5', 'get_folder_file_list__9aea5fe5a12921799c11d30552224cd9', 'get_pdf_page_and_orientation__718b8e151527c1d36352c5376a93ead7', 'get_csv_in_documents__a95ba08342a252657d47056ba988912e', 'get_folder_and_files__2cb8499f32af5ada0bb0e955434e59dd', 'get_file_rename_status__5128f87c595d3e5765ba37b7dff4c652', 'get_zip_file_list__b6394280513f3b0b8d3eaa6ea7ea0992', 'get_local_folder_check__1d547c7422f60e4be92e89ee2379630a', 'get_file_content__d8b2f0ccd335f9a3a3203363805b5ac0', 'get_doc_file_list__d6bd2cbacad8b9b51e1e935f00211e4d', 'get_file_exists__0ae8da7e0d709c03ee846a1d011df875', 'get_text_file_content__ec3e3cd160378737f8c3074a7c1fe7f0', 'get_dir_move_state__a7ceb843bee94531ff589f6c5602590a', 'get_import_file__c191ddd519b4043306a4effe99ae0aac', 'get_dir_structure__f789180956e4b0547b5929aa18a0352e', 'get_txt_file_content__f7a75c940bb0ba7687c094e66be96f16', 'get_text_file_content__55e51bced75ba3e5e08e1be190df564e', 'get_text_file_content__358385204fa7b84b2d9118b3314be461', 'get_txt_file_lines__5e14eceba731d7ebe21b244525a2edc1', 'get_dir_ownership__425163a454302c61dc411d05ac9420b7', 'get_text_file_content__ae954eaf7f0ec05c9c2918c4eec73bcc', 'get_total_file_content__aa615bdde918e746d9793d851a2bb4f8', 'get_folder_file_list__1a224fb0890014daa398af667e7e35c0', 'get_file_system_check__f073ea69e2970c4ed8abb5029189b0ac', 'get_vm_file_content__28ae802c5923a6e7ecfebd0e5cc59702', 'get_sorted_invoice_folders__5d7cb95eaab32c834b0336c7360277e1', 'get_dir_structure__95c704f58705a26de50af292b2eb5059', 'get_manifest_json__61a5e99298254bc3ea7017f23352ee68', 'get_file_info__fda9115554280abf401514b41e43c7b7', 'get_frame_extract_files__5e444f1b525a6375eead73fcf3fce163', 'get_pdf_file_check__1c15db0bd188a0b79dfb455b9076c68e', 'get_vm_file_content__a8cfa9828c8ad4ee330aec94e4adaf4c', 'get_docx_text_content__5cd9f8324e43d00ab6826873ca723208', 'get_lecture_pdf_exists__29ea2d20dd9405b64f3df8f5049b79d3', 'get_stats_file__7636602600807a890f06c684e5901e16', 'get_docx_line_content__31146efdc87a4defa3f97309b7095c90', 'get_pdf_page_info__075106d955ead60b31af19408ecd54ac', 'get_file_content__a4d1b434aaedff8ce21d5e5d82eac67e', 'get_file_content__6790a81b7f7c76a950ade3c3fcc6cc45', 'get_manifest_json__783369d0e7e20c29051ab5abdc75ebb3', 'get_pdf_on_desktop__18b0663c1e76805fd4ea486cbf72380b', 'get_file_content_check__2372f00ad964142d3685a78ced2bb21c', 'get_pdf_file_info__beb4b25746cd97a98dc95327476aaf4d', 'get_vm_files_multi__8aa2e8e2be1ca5d610e977b9c78060b2', 'get_result_file_content__4c5e4ba8ba69ff54d1e2443439646d44', 'get_csv_head__8b9d6ae5a27e52886aef51e39cc4df8a', 'get_csv_content__f43838d5904be041721c003e00fd36b5', 'get_check_output_file__f4c7cbbe6d6bd73cc0b521d770fb3945', 'get_file_info__5d2d411fc8203609cb2ed6eef9424f86', 'get_result_file_content__4128b5119649f756a2c7d9fc4321d25f', 'get_merged_novel_file__2abe2760620ddfe22faad3f384e1cae4', 'get_local_folder_list__91f316cac3c13e1ffa3f1d0a1879bf2b', 'get_desktop_file__bc88c1a76af0c575a8c22634fe96ef8b_qw35sft2_f69c4fb7', 'get_downloads_pdf__5416755d2dd223fd56c07a64ea26e507_qw35sft2_5814a28c', 'get_file_exists__de94d61f074111bc63d9a066e03aa46f_qw35sft2_73b15751', 'get_csv_and_pdf__f45249b0314207a95d7366d15bee2907_qw35sft2_1e85c06e', 'get_bold_header_and_pdf__86ea3ac794566dae35e2163b9c7d3c94_qw35sft2_e37cd2f8', 'get_col_a_bold_and_csv__8f6a3ccc5fb938975de38e8fbcc8f580_qw35sft2_725bcefe', 'get_file_exists_on_desktop__3ed023e79bfe94e32e5ea286856ad04c_qw35sft2_b66cefae', 'get_pdf_exists__19cbc30a1547517beef14a189cce767e_qw35sft2_652518fc', 'get_pdf_and_odt_exists__c65906d6f690144adc013e7fcea2301e_qw35sft2_5e4aefbc', 'get_para0_word_content__469a4ad5eae55be7063b8ec35b77b37d_qw35sft2_947cb49d', 'get_orgsummary_pdf_exists__82b855263e1bdeb6cb9c2f640e53bf6c_qw35sft2_1748291e', 'get_pdf_in_documents__df1f48650ab6f100a1208b8040cd8828_qw35sft2_6e48e714', 'get_pdf_and_footer_pagenum__0bb862e1956085267ced241046059b10_qw35sft2_f87e83a2', 'get_path_text_file__a43735a3f0450074dfa18129f92a73a9_qw35sft2_74219b39', 'get_text_file_content__f36e167c6b956664e01e5e47712b671b_qw35sft2_c7971e50', 'get_file_content__5ca68902218b7e30665245f17696931d_qw35sft2_6875fbcc', 'get_desktop_file_exists__6964660b6570aa960a367d963aeb2477_qw35sft2_50044e66', 'get_dir_exists__e66cc6b0ec6479dc303b558a9cc72642_qw35sft2_e41816cf', 'get_res_txt_content__da27b9fa71ba67953d0b9f649a3197fa_qw35sft2_ead3d6e3', 'get_receipts_file_listing__d248703d35223088b5508b0c05332ea3_qw35sft2_cbb3b1f6', 'get_text_file_content__30cadb0ec7eca28df693cf924ab88301_qw35sft2_a743245c', 'get_fs_file_exists__21cbb74de44bf7839b1ea14bdfabea03_qw35sft2_9687eb5c', 'get_txt_multihop_gemini__61501dfdaa17fec1d3c116cb7f744d5c_qw35sft2_1fcc099e', 'get_csv_filtered_a__2fadb675c56bf8de284fee43324487b2_qw35sft2_b0731958', 'get_dir_exists__e103d8977cd2fefeec7972ff8169e36d_qw35sft2_0dfba723', 'get_file_at_path__82c707062135fc7c567ff519e1b5a575_qw35sft2_273ea9cf', 'get_bubblesort_both_files__8c65eee451a1bb106a548123af80ea08_qw35sft2_60135b55', 'get_main_py_content__f65942d684aa73fa8a727494a6335877_qw35sft2_9e5513a5', 'get_csv_space_merged__00d80996b3209df110a0b1e69fc0ab31_qw35sft2_b3c91bc9', 'get_file_exists__057715c74cd9cb2c93d92377c04bb077_qw35sft2_722a620b', 'get_clipboard_content__241466d2f7566981ed09035f56716bb4_qw35sft2_9c0b84e6', 'get_res_txt_content__68c8947721f1221211aba42cd6d1735d_qw35sft2_37d291c9', 'get_file_state__ea8c7a7a44f5cb78e25fd89d658a863f_qw35sft2_fc8fdbc1', 'get_all_file_permissions__1540d35f6307434ff998913b82a97dfa_qw35sft2_ae6ad5fe', 'get_user_identity_file__9be4c7481abd67d84bd529b63ed66ca8_qw35sft2_9dcd61e2', 'get_count_file__5799d2a2e180d6540c452f203d2bde46_qw35sft2_55233ea8', 'get_copy_4dirs__4f3015d5aa890d8fea505e363d3f7aff_qw35sft2_2fd46445', 'get_file_perms_and_rename__653ac9883bb8d908f613f5a6b4d6a405_qw35sft2_68e3d2be', 'get_sys_info_file__5cefed7a5e164dc00f5c88ae5726c23e_qw35sft2_5a973c8c', 'get_old_files_cleanup__aba3d770f03c5889eb920c2ccea4f4e0_qw35sft2_1335d4ce', 'get_copy_and_create_file2__0f2ab16ddcc8f1a930b6b27f18088a10_qw35sft2_1fd84deb', 'get_file_and_dir_permissions__6736312e2a9f6c78caa7e5e2f35a4133_qw35sft2_e7b0794b', 'get_user_audit_files__c74bb543554ec2beed49099c34c9413f_qw35sft2_d16d2df5', 'get_file_org_state__8f46312161f8fe4e75948b31954b5be0_qw35sft2_c656985f', 'get_subdir_split_permissions__31b8299b7936d108e594f17a62518506_qw35sft2_a0f0f736', 'get_file_copy_3dirs__d110c9f1795d3fe99db0ae18f13e5b66_qw35sft2_3a1d9288', 'get_desktop_files_check__ee992b5c4f2de4a7b0ef0f22ec1a67b1_qw35sft2_e9092f55', 'get_rename_and_subfolder__cc8a90beff8b6d7aa9d9f79938c158e4_qw35sft2_42301119', 'get_new_files_perms__2e6f4a565bdf8ce8ae9bb4e56e1c79f7_qw35sft2_218fed0a', 'get_home_users_file__baeeb37009528a1d856ed9e685fd5d89_qw35sft2_de61165a', 'get_two_file_state__bb3776631f95d3c4e2bcd96582b451ff_qw35sft2_a5bedd98', 'get_rename_and_file__caa3a14501f9d1aa7da658fbd50adb72_qw35sft2_ca1c5ca0', 'get_file_perms_and_archive__00b9eddcbcb085dc89e5e17c3b839024_qw35sft2_9bf6c2cc', 'get_archive_content__7b7da38700f679a895053728c5a7d35b_qw35sft2_1ee11bd6', 'get_two_dir_notebook_split__e71563318b46b13cadfcb76e6b9ce246_qw35sft2_71e98bf9', 'get_bills_and_local_folders__0b7b1b3c91e02a8222608f02952fc214_qw35sft2_8ba964c8', 'get_file_exists__df188ad33aeda3e315e62cc2c6afb173_qw35sft2_f28ae247', 'get_main_py_content__2c85ba49c59e818794c6aa64361a93f2_qw35sft2_47dcf891', 'get_main_py_content__4c84b586b63f7c1cf9b6df39aedcd5ed_qw35sft2_db7ad539']

def get_file_dual_location__454f85dd5a1086d0f2418681bdac8477(env, config: dict):
    """Check if file exists at two locations (original restored + copy)."""
    path_a = config.get('path_a', '')
    path_b = config.get('path_b', '')
    check_cmd = f"[ -f '{path_a}' ] && echo 'a_exists' || echo 'a_missing'; [ -f '{path_b}' ] && echo 'b_exists' || echo 'b_missing'"
    result = env.controller.run_bash_script(check_cmd, timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    return {'location_a': 'a_exists' in output, 'location_b': 'b_exists' in output}

def get_manifest_json__ab85f9bbe05040b98eb1a964cbee7866(env, config: dict):
    """Read and parse manifest.json from the VM."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {'manifest': data}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_odt_files_exist__016af9433419594b3dc790848ad51e3b(env, config: dict):
    """Check that .odt files exist on Desktop and were converted from .doc files."""
    vm_ip = env.vm_ip
    port = env.server_port
    history_cmd = ['/bin/bash', '-c', 'output=$(cat ~/.bash_history | grep -E "(soffice|libreoffice).+--convert-to\\s+odt.+\\*\\.doc"); if [ -z "$output" ]; then echo "no_command"; else echo "command_found"; fi']
    resp_hist = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': history_cmd, 'shell': False})
    history_result = ''
    if resp_hist.status_code == 200:
        history_result = resp_hist.json().get('output', '').strip()
    count_cmd = ['/bin/bash', '-c', 'ls /home/user/Desktop/*.odt 2>/dev/null | wc -l']
    resp_count = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': count_cmd, 'shell': False})
    odt_count = 0
    if resp_count.status_code == 200:
        try:
            odt_count = int(resp_count.json().get('output', '0').strip())
        except ValueError:
            odt_count = 0
    return {'command_found': history_result == 'command_found', 'odt_count': odt_count}

def get_text_file_content__b150ff4c0a473c93c912a80df923233a(env, config: dict):
    """Read a text file from VM and return its content as a list of stripped lines."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'lines': []}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        return {'lines': lines, 'raw': content}
    except Exception as e:
        return {'error': str(e), 'lines': []}

def get_folder_file_list__448eb54c6f660c787448d046dde3cea3(env, config: dict):
    """List files in a specified directory on the VM."""
    folder_path = config.get('path', '')
    result = env.controller.run_bash_script(f'ls -1 "{folder_path}" 2>/dev/null | sort', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if not output:
        return {'files': [], 'count': 0, 'exists': False}
    files = [f for f in output.split('\n') if f.strip()]
    return {'files': sorted(files), 'count': len(files), 'exists': True}

def get_file_move_check__e60640f858dbfcce19e818f08302a608(env, config: dict):
    """Check if a file exists at destination and not at source."""
    file_name = config.get('file_name', '')
    src_dir = config.get('src_dir', '/home/user/Pictures')
    dst_dir = config.get('dst_dir', '/home/user/Desktop')
    src_check = env.controller.run_bash_script(f"test -f {src_dir}/{file_name} && echo 'EXISTS' || echo 'MISSING'", timeout=30)
    src_exists = 'EXISTS' in (src_check.get('output', '') if isinstance(src_check, dict) else str(src_check))
    dst_check = env.controller.run_bash_script(f"test -f {dst_dir}/{file_name} && echo 'EXISTS' || echo 'MISSING'", timeout=30)
    dst_exists = 'EXISTS' in (dst_check.get('output', '') if isinstance(dst_check, dict) else str(dst_check))
    return {'src_exists': src_exists, 'dst_exists': dst_exists}

def get_pdf_page_info__a8da1aa98de37fa62b8d321ad0589eaf(env, config: dict):
    """Get PDF page count from a file on the VM."""
    import tempfile, os
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        try:
            from pypdf import PdfReader
        except ImportError:
            return {'error': 'PyPDF2/pypdf not available'}
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        reader = PdfReader(tmp_path)
        num_pages = len(reader.pages)
        return {'num_pages': num_pages, 'file_exists': True}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_main_py_content__bef206759cc52e3fdba71a4f0e77e8b3(env, config: dict):
    """Read main.py content from VM."""
    file_path = config.get('path', '/home/user/project/main.py')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        content = file_bytes.decode('utf-8')
        return {'content': content}
    except Exception as e:
        return {'error': str(e)}

def get_txt_file_lines__d445ea0218043b2b599f555fb0b0c923(env, config: dict):
    """Read a text file from the VM and return its non-empty lines."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'lines': []}
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return {'lines': lines, 'count': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_text_file_content__ca2a0f6f5509af90b7b3cf524daf3368(env, config: dict):
    """Get content of a text file from VM."""
    file_path = config.get('path', '/home/user/Desktop/grf23.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
        return {'content': content, 'length': len(content)}
    except Exception as e:
        return {'error': str(e), 'content': ''}

def get_file_content__4c4a9bd632c3f7487e991bbd8a3f942a(env, config: dict):
    """Get text file content from VM."""
    path = config.get('path', '')
    try:
        result = env.controller.run_bash_script(f'cat "{path}" 2>/dev/null', timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        if not output:
            return {'error': 'File not found or empty', 'exists': False, 'content': ''}
        return {'exists': True, 'content': output}
    except Exception as e:
        return {'error': str(e), 'exists': False, 'content': ''}

def get_verified_folder_files__445a1e82d5445803b67308e02fdc3544(env, config: dict):
    """List PDF files in the verified folder on the Desktop."""
    folder_path = config.get('folder_path', '/home/user/Desktop/verified')
    result = env.controller.run_bash_script(f'ls -1 "{folder_path}" 2>/dev/null', timeout=30)
    stdout = result.get('output', '') if isinstance(result, dict) else str(result)
    files = [f.strip() for f in stdout.strip().split('\n') if f.strip()]
    return {'files': files, 'count': len(files)}

def get_docx_to_pdf_state__8f355d6af58b179cf68cc37725d92b0c(env, config: dict):
    """Check that .docx files were converted to PDF via command line."""
    vm_ip = env.vm_ip
    port = env.server_port
    history_cmd = ['/bin/bash', '-c', 'output=$(cat ~/.bash_history | grep -E "(soffice|libreoffice).+--convert-to\\s+pdf.+\\*\\.docx"); if [ -z "$output" ]; then echo "no_command"; else echo "command_found"; fi']
    resp_hist = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': history_cmd, 'shell': False})
    history_result = ''
    if resp_hist.status_code == 200:
        history_result = resp_hist.json().get('output', '').strip()
    check_cmd = ['/bin/bash', '-c', 'count=0; for f in /home/user/Desktop/1cfc37e1-344f-52ef-be9f-69bc1855316d.pdf /home/user/Desktop/adfddc91-944f-572e-8439-c67b67f4c5f7.pdf /home/user/Desktop/d119ec81-bc60-5b76-8679-5f942b191c70.pdf; do if [ -f "$f" ]; then count=$((count+1)); fi; done; echo $count']
    resp_count = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': check_cmd, 'shell': False})
    pdf_count = 0
    if resp_count.status_code == 200:
        try:
            pdf_count = int(resp_count.json().get('output', '0').strip())
        except ValueError:
            pdf_count = 0
    return {'command_found': history_result == 'command_found', 'docx_pdf_count': pdf_count}

def get_pdf_page_count__ddacf8dd85e44a11a4923ff07ce0abff(env, config: dict):
    """Get PDF file existence and page count from VM."""
    import tempfile
    import os
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False, 'page_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        return {'exists': True, 'page_count': page_count}
    except Exception as e:
        return {'error': str(e), 'exists': False, 'page_count': 0}
    finally:
        os.unlink(tmp_path)

def get_file_content__29f89216bfdb55e864e9190350cbf19a(env, config: dict):
    """Read a file from the VM and return its content and size."""
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'size': 0, 'content': ''}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
    except Exception:
        content = ''
    return {'exists': True, 'size': len(file_bytes), 'content': content}

def get_csv_file_check__02188a8f533477bfdcad341f72281666(env, config: dict):
    """Read CSV file from VM and return its structure info."""
    file_path = config.get('path', '/home/user/Desktop/stock_data.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False, mode='wb') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            reader = csv.reader(f)
            rows = list(reader)
        if len(rows) == 0:
            return {'error': 'Empty CSV file', 'exists': True, 'row_count': 0}
        header = rows[0]
        data_rows = rows[1:]
        non_empty_count = 0
        for row in data_rows:
            if len(row) > 0 and row[0].strip():
                non_empty_count += 1
        return {'exists': True, 'header': header, 'header_count': len(header), 'total_rows': len(rows), 'data_rows': non_empty_count}
    finally:
        os.unlink(tmp_path)

def get_user_homedir__df678257c48d0ebe8f2e03968592b783(env, config: dict):
    """Get user home directory and existence status from the system."""
    username = config.get('username', 'charles')
    try:
        result = env.controller.run_bash_script(f'getent passwd {username}', timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        if not output:
            return {'exists': False, 'home': '', 'shell': ''}
        parts = output.split(':')
        if len(parts) >= 7:
            return {'exists': True, 'home': parts[5], 'shell': parts[6]}
        return {'exists': False, 'home': '', 'shell': ''}
    except Exception as e:
        logger.error(f'Error getting user info: {e}')
        return {'exists': False, 'home': '', 'shell': ''}

def get_text_file_content__204b2db58c306983ce47277dbc39788e(env, config: dict):
    """Read a text file from the VM and return its stripped content."""
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File not found: {path}')
            return {'error': 'File not found', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return {'error': str(e), 'content': ''}

def get_file_organize_state__a65eb6fd760478532b3e118be8d079f9(env, config: dict):
    """Check if receipts directory exists and OIP.jpg is inside it."""
    try:
        dir_check = env.controller.run_bash_script('test -d /home/user/receipts && echo "DIR_EXISTS" || echo "DIR_MISSING"', timeout=30)
        dir_exists = 'DIR_EXISTS' in dir_check.get('output', '')
        file_check = env.controller.run_bash_script('test -f /home/user/receipts/OIP.jpg && echo "FILE_EXISTS" || echo "FILE_MISSING"', timeout=30)
        file_exists = 'FILE_EXISTS' in file_check.get('output', '')
        return {'dir_exists': dir_exists, 'file_in_dir': file_exists}
    except Exception as e:
        return {'error': str(e)}

def get_file_copy_and_clipboard__267da608b7e8dd72021770568c2da84f(env, config: dict):
    """Get file existence at target path and clipboard content from VM."""
    vm_ip = env.vm_ip
    port = env.server_port
    target_path = config.get('target_path', '')
    file_check_cmd = f"test -f '{target_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'"
    try:
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': file_check_cmd, 'shell': True})
        file_exists = False
        if response.status_code == 200:
            file_exists = 'EXISTS' in response.json().get('output', '')
    except Exception as e:
        logger.error('Failed to check file existence: %s', e)
        file_exists = False
    clip_cmd = 'xsel --clipboard --output'
    try:
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': clip_cmd, 'shell': True})
        clipboard = ''
        if response.status_code == 200:
            clipboard = response.json().get('output', '').strip()
    except Exception as e:
        logger.error('Failed to get clipboard content: %s', e)
        clipboard = ''
    return {'file_exists': file_exists, 'clipboard': clipboard}

def get_bib_file_content__dd385e64b1e5f06bf4b4c074020cb6a0(env, config: dict):
    """Get content of references.bib file from VM."""
    try:
        file_path = config.get('path', '/home/user/Desktop/references.bib')
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'content': '', 'error': 'File not found or empty'}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        logger.error(f'Error reading bib file: {e}')
        return {'content': '', 'error': str(e)}

def get_song_list_file__09209b5893b5a25bd669841ecee9c784(env, config: dict):
    """Read content of ~/Music/song_list.txt from VM."""
    try:
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Music/song_list.txt'))
        if not file_bytes:
            return {'error': 'File not found', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        return {'error': str(e), 'content': ''}

def get_docx_text_content__b98b6fc5ce8ebe003cf7eaeacc71d919(env, config: dict):
    """Get all text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'text': '', 'paragraphs': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = '\n'.join(paragraphs)
        return {'text': full_text, 'paragraphs': paragraphs, 'paragraph_count': len(paragraphs)}
    finally:
        os.unlink(tmp_path)

def get_file_rename_check__9b65526186f9fb77c93ff9d7de5c5fcb(env, config: dict):
    """Check files in Pictures directory to verify renaming."""
    result = env.controller.run_bash_script("ls /home/user/Pictures/ 2>/dev/null || echo '__ERROR__'", timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    if '__ERROR__' in output:
        return {'error': 'Cannot list Pictures directory'}
    files = [f.strip() for f in output.strip().split('\n') if f.strip()]
    return {'files': files}

def get_def_file__9f211d2fef2982f0959041bc806f53c4(env, config: dict):
    """Read the function/class definitions file from VM and return parsed content."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = [l.strip() for l in content.strip().split('\n') if l.strip()]
    def_lines = [l for l in lines if l.startswith('def ') or l.startswith('class ')]
    non_def_lines = [l for l in lines if not (l.startswith('def ') or l.startswith('class '))]
    return {'content': content, 'lines': lines, 'line_count': len(lines), 'def_lines': def_lines, 'non_def_lines': non_def_lines, 'def_count': len(def_lines)}

def get_dir_existence__93ffa0d39e3e07dab53be514a31da362(env, config: dict):
    """Check if a directory exists on the VM by running a test command."""
    dir_path = config.get('dir_path', '')
    result = env.controller.run_bash_script(f'test -d "{dir_path}" && echo "EXISTS" || echo "NOT_EXISTS"', timeout=10)
    if isinstance(result, dict):
        output = result.get('output', '').strip()
    else:
        output = str(result).strip() if result else ''
    return {'exists': 'EXISTS' in output, 'dir_path': dir_path}

def get_vm_file_content__d175f88bb1991b0561f10b03c2a108ce(env, config: dict):
    """Get file content from VM for iostat disk I/O report."""
    file_path = config.get('path', '')
    dest = config.get('dest', 'result.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        cache_dir = env.cache_dir if hasattr(env, 'cache_dir') else '/tmp'
        dest_path = os.path.join(cache_dir, dest)
        with open(dest_path, 'wb') as f:
            f.write(file_bytes)
        with open(dest_path, 'r') as f:
            content = f.read()
        return {'content': content, 'file_path': dest_path}
    except Exception as e:
        logger.error(f'get_vm_file_content error: {e}')
        return {'error': str(e)}

def get_vm_file_content__6c3110d7dd0b28f9da7376f5ce78a46b(env, config: dict):
    """Get file content from VM as a string."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
    except Exception:
        content = ''
    return {'content': content}

def get_music_folder_state__eed664169e1770df941a91e7673f0e98(env, config: dict):
    """Check directory existence and file presence in ~/Music/Classics/."""
    try:
        dir_check = env.controller.run_bash_script('test -d /home/user/Music/Classics && echo "EXISTS" || echo "NOT_EXISTS"', timeout=30)
        dir_exists = 'EXISTS' in dir_check.get('output', '')
        files_in_classics = []
        if dir_exists:
            ls_result = env.controller.run_bash_script('ls /home/user/Music/Classics/', timeout=30)
            files_in_classics = [f.strip() for f in ls_result.get('output', '').strip().split('\n') if f.strip()]
        ls_music = env.controller.run_bash_script('ls /home/user/Music/*.mp3 2>/dev/null || true', timeout=30)
        files_in_music = [f.strip().split('/')[-1] for f in ls_music.get('output', '').strip().split('\n') if f.strip()]
        return {'dir_exists': dir_exists, 'files_in_classics': files_in_classics, 'files_in_music_root': files_in_music}
    except Exception as e:
        return {'error': str(e)}

def get_targz_contents__d544b10f933e46c3b8f208241fda6c7e(env, config: dict):
    """Download tar.gz from VM and list its contents."""
    import tempfile
    import os
    import tarfile
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'exists': False, 'files': [], 'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.tar.gz', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with tarfile.open(tmp_path, 'r:gz') as tar:
                members = tar.getnames()
                basenames = [os.path.basename(m) for m in members if os.path.basename(m)]
            return {'exists': True, 'files': basenames, 'raw_members': members}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'files': [], 'error': str(e)}

def get_pdf_list_file__f72fa74134cf5a58d1281947ac1e62bd(env, config: dict):
    """Get the PDF list file from Desktop and return its content."""
    file_path = config.get('path', '/home/user/Desktop/pdf_list.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'exists': False}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'exists': True, 'content': content, 'content_lower': content.lower()}
    except Exception as e:
        logger.error(f'Error reading PDF list file: {e}')
        return {'error': str(e), 'exists': False}

def get_file_content__79ff096ee36ebe0503f424e88a332ca2(env, config: dict):
    """Read text file content from VM."""
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        content = file_bytes.decode('utf-8').strip()
        return {'content': content}
    except Exception as e:
        return {'error': str(e)}

def get_docx_text_content__5ac00db9c05416b4087c6a9ef6d0dee7(env, config: dict):
    """Get all text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'text': '', 'paragraphs': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = '\n'.join(paragraphs)
        return {'text': full_text, 'paragraphs': paragraphs, 'paragraph_count': len(paragraphs)}
    finally:
        os.unlink(tmp_path)

def get_docx_content__2aeb38ddfad7ab2f7bd9e6f63f6cc1bf(env, config: dict):
    """Get text content and image count from a .docx file on the VM."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'text': '', 'image_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        text = '\n'.join([p.text for p in doc.paragraphs])
        image_count = 0
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            image_files = [f for f in zf.namelist() if f.startswith('word/media/')]
            image_count = len(image_files)
        return {'text': text, 'image_count': image_count}
    except Exception as e:
        logger.error(f'Error reading docx: {e}')
        return {'error': str(e), 'text': '', 'image_count': 0}
    finally:
        os.unlink(tmp_path)

def get_ods_to_pdf_check__40ad56f2b3098a6972e48b06b58c95e0(env, config: dict):
    """Check if ODS was converted to PDF via terminal command."""
    result = {}
    try:
        history_output = env.controller.run_bash_script("cat ~/.bash_history | grep '\\(soffice\\|libreoffice\\).*--convert-to\\s\\+pdf'", timeout=30)
        history_text = history_output.get('output', '') if isinstance(history_output, dict) else str(history_output)
        result['used_terminal'] = 'use terminal' if history_text.strip() else 'use no terminal'
    except Exception:
        result['used_terminal'] = 'use no terminal'
    pdf_path = config.get('path', '/home/user/Desktop/file_example_ODS_5000.pdf')
    try:
        file_bytes = env.controller.get_file(pdf_path)
        if file_bytes and len(file_bytes) > 0:
            result['file_exists'] = True
            result['file_size'] = len(file_bytes)
            result['is_valid_pdf'] = file_bytes[:5] == b'%PDF-'
        else:
            result['file_exists'] = False
            result['file_size'] = 0
            result['is_valid_pdf'] = False
    except Exception:
        result['file_exists'] = False
        result['file_size'] = 0
        result['is_valid_pdf'] = False
    return result

def get_folder_file_list__69659eabd73459db19e2b0d1637cc759(env, config: dict):
    """List files in the specified folder on the VM."""
    folder_path = config.get('path', '')
    result = env.controller.run_bash_script(f'ls "{folder_path}" 2>/dev/null', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if not output:
        return {'files': [], 'raw': ''}
    files = sorted([f for f in output.split('\n') if f.strip()])
    return {'files': files, 'raw': output}

def get_novel_info_file__6a9459e3699067e60bf00ce031fcbd18(env, config: dict):
    """Get the novel info file from Desktop and return its content."""
    file_path = config.get('path', '/home/user/Desktop/novel_info.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'exists': False}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'exists': True, 'content': content, 'content_lower': content.lower()}
    except Exception as e:
        logger.error(f'Error reading novel info file: {e}')
        return {'error': str(e), 'exists': False}

def get_bib_file_content__db9975908e0c2c40a4467f2add6ed446(env, config: dict):
    """Get content of references.bib file from VM."""
    try:
        file_path = config.get('path', '/home/user/Desktop/references.bib')
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'content': '', 'error': 'File not found or empty'}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        logger.error(f'Error reading bib file: {e}')
        return {'content': '', 'error': str(e)}

def get_vm_file(env, config: Dict[str, Any]):
    """Get file from VM - simplified version for this getter."""
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        return cache_path
    except Exception as e:
        logger.error(f'Error getting file {path}: {e}')
        return None

def get_file_content__05446d05a2730db5bc45314af39c65c5(env, config: dict):
    """Read text file content from VM."""
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8').strip()
        return {'content': content}
    except Exception as e:
        return {'error': str(e), 'content': ''}

def get_folder_file_list__9aea5fe5a12921799c11d30552224cd9(env, config: dict):
    """List files in the specified folder on the VM."""
    folder_path = config.get('path', '')
    result = env.controller.run_bash_script(f'ls "{folder_path}" 2>/dev/null', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if not output:
        return {'files': [], 'raw': ''}
    files = sorted([f for f in output.split('\n') if f.strip()])
    return {'files': files, 'raw': output}

def get_pdf_page_and_orientation__718b8e151527c1d36352c5376a93ead7(env, config: dict):
    """Get PDF page count and orientation (landscape vs portrait)."""
    import tempfile
    import os
    file_path = config.get('path', '/home/user/Resize_Cells_Fit_Page.pdf')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        if page_count > 0:
            page = reader.pages[0]
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            is_landscape = width > height
        else:
            width = 0
            height = 0
            is_landscape = False
        return {'page_count': page_count, 'is_landscape': is_landscape, 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_csv_in_documents__a95ba08342a252657d47056ba988912e(env, config: dict):
    """Check if ODS was converted to CSV and saved to Documents directory."""
    result = {}
    try:
        history_output = env.controller.run_bash_script("cat ~/.bash_history | grep '\\(soffice\\|libreoffice\\).*--convert-to\\s\\+csv'", timeout=30)
        history_text = history_output.get('output', '') if isinstance(history_output, dict) else str(history_output)
        result['used_terminal'] = 'use terminal' if history_text.strip() else 'use no terminal'
    except Exception:
        result['used_terminal'] = 'use no terminal'
    csv_path = config.get('path', '/home/user/Documents/file_example_ODS_5000.csv')
    try:
        file_bytes = env.controller.get_file(csv_path)
        if file_bytes and len(file_bytes) > 0:
            content = file_bytes.decode('utf-8', errors='replace')
            lines = content.strip().split('\n')
            result['file_exists'] = True
            result['line_count'] = len(lines)
            header = lines[0] if lines else ''
            result['has_valid_header'] = 'First Name' in header and 'Last Name' in header
        else:
            result['file_exists'] = False
            result['line_count'] = 0
            result['has_valid_header'] = False
    except Exception:
        result['file_exists'] = False
        result['line_count'] = 0
        result['has_valid_header'] = False
    return result

def get_folder_and_files__2cb8499f32af5ada0bb0e955434e59dd(env, config: dict):
    """Check if folder exists and what files are in Pictures and subfolder."""
    folder_name = config.get('folder_name', 'Mountains')
    base_dir = '/home/user/Pictures'
    folder_check = env.controller.run_bash_script(f"test -d {base_dir}/{folder_name} && echo 'EXISTS' || echo 'MISSING'", timeout=30)
    folder_exists = 'EXISTS' in (folder_check.get('output', '') if isinstance(folder_check, dict) else str(folder_check))
    subfolder_result = env.controller.run_bash_script(f"ls {base_dir}/{folder_name}/ 2>/dev/null || echo ''", timeout=30)
    subfolder_output = subfolder_result.get('output', '') if isinstance(subfolder_result, dict) else str(subfolder_result)
    subfolder_files = [f.strip() for f in subfolder_output.strip().split('\n') if f.strip()]
    root_result = env.controller.run_bash_script(f'ls -p {base_dir}/ 2>/dev/null | grep -v /', timeout=30)
    root_output = root_result.get('output', '') if isinstance(root_result, dict) else str(root_result)
    root_files = [f.strip() for f in root_output.strip().split('\n') if f.strip()]
    return {'folder_exists': folder_exists, 'subfolder_files': subfolder_files, 'root_files': root_files}

def get_file_rename_status__5128f87c595d3e5765ba37b7dff4c652(env, config: dict):
    """Check if file was renamed correctly by verifying new name exists and old name is gone."""
    new_path = config.get('new_path', '')
    old_path = config.get('old_path', '')
    check_cmd = f"[ -f '{new_path}' ] && echo 'new_exists' || echo 'new_missing'; [ -f '{old_path}' ] && echo 'old_exists' || echo 'old_missing'"
    result = env.controller.run_bash_script(check_cmd, timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    return {'new_exists': 'new_exists' in output, 'old_gone': 'old_missing' in output}

def get_zip_file_list__b6394280513f3b0b8d3eaa6ea7ea0992(env, config: dict):
    """List files inside a zip archive on the VM."""
    zip_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'test -f "{zip_path}" && echo "exists" || echo "missing"', timeout=30)
    exists_output = exists_result.get('output', '').strip() if isinstance(exists_result, dict) else ''
    if 'missing' in exists_output or not exists_output:
        return {'files': [], 'count': 0, 'exists': False}
    result = env.controller.run_bash_script(f'''unzip -l "{zip_path}" 2>/dev/null | tail -n +4 | head -n -2 | awk '{{print $NF}}' | sort''', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if not output:
        return {'files': [], 'count': 0, 'exists': True}
    files = [f.split('/')[-1] for f in output.split('\n') if f.strip() and (not f.strip().endswith('/'))]
    return {'files': sorted(files), 'count': len(files), 'exists': True}

def get_local_folder_check__1d547c7422f60e4be92e89ee2379630a(env, config: dict):
    """Check if a specific folder exists under Local Folders in Thunderbird."""
    folder_name = config.get('folder_name', 'Projects')
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    local_folders_path = f'{profile_path}/Mail/Local Folders'
    result = env.controller.run_bash_script(f"ls -1 '{local_folders_path}/'", timeout=30)
    if isinstance(result, dict):
        output = result.get('output', '')
    else:
        output = str(result)
    files = [f.strip() for f in output.strip().split('\n') if f.strip()]
    folder_exists = folder_name in files or f'{folder_name}.msf' in files
    return {'folder_exists': folder_exists, 'folder_name': folder_name, 'files_found': files}

def get_file_content__d8b2f0ccd335f9a3a3203363805b5ac0(env, config: dict):
    """Download file from VM and return its content."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace')
    return {'content': content}

def get_doc_file_list__d6bd2cbacad8b9b51e1e935f00211e4d(env, config: dict):
    """Get the content of filelist.txt and count of actual .doc files."""
    vm_ip = env.vm_ip
    port = env.server_port
    filelist_cmd = ['/bin/bash', '-c', "cat /home/user/Desktop/filelist.txt 2>/dev/null || echo '__FILE_NOT_FOUND__'"]
    resp_list = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': filelist_cmd, 'shell': False})
    filelist_content = ''
    if resp_list.status_code == 200:
        filelist_content = resp_list.json().get('output', '').strip()
    count_cmd = ['/bin/bash', '-c', 'ls /home/user/Desktop/*.doc 2>/dev/null | wc -l']
    resp_count = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': count_cmd, 'shell': False})
    actual_doc_count = 0
    if resp_count.status_code == 200:
        try:
            actual_doc_count = int(resp_count.json().get('output', '0').strip())
        except ValueError:
            actual_doc_count = 0
    return {'filelist_content': filelist_content, 'filelist_exists': '__FILE_NOT_FOUND__' not in filelist_content, 'actual_doc_count': actual_doc_count}

def get_file_exists__0ae8da7e0d709c03ee846a1d011df875(env, config: dict):
    """Check if a file exists at the specified path on the VM."""
    target_path = config.get('path', '')
    result = env.controller.run_bash_script(f'test -f "{target_path}" && echo "EXISTS" || echo "NOT_FOUND"', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    file_size_result = env.controller.run_bash_script(f'stat -c %s "{target_path}" 2>/dev/null || echo "0"', timeout=30)
    size_output = file_size_result.get('output', '0').strip() if isinstance(file_size_result, dict) else str(file_size_result).strip()
    try:
        file_size = int(size_output)
    except ValueError:
        file_size = 0
    return {'exists': 'EXISTS' in output, 'file_size': file_size, 'path': target_path}

def get_text_file_content__ec3e3cd160378737f8c3074a7c1fe7f0(env, config: dict):
    """Read text file content from VM."""
    target_path = config.get('path', '')
    file_bytes = env.controller.get_file(target_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8').strip()
    except (UnicodeDecodeError, AttributeError):
        return {'error': 'Cannot decode file', 'content': ''}
    return {'content': content, 'path': target_path}

def get_dir_move_state__a7ceb843bee94531ff589f6c5602590a(env, config: dict):
    """Check directory state after move operation - whether dir4 moved from dir3 to dir1."""
    results = {}
    try:
        r1 = env.controller.run_bash_script('test -d /home/user/Desktop/dir1/dir4 && echo YES || echo NO', timeout=30)
        output1 = r1.get('output', r1.get('stdout', '')).strip() if isinstance(r1, dict) else str(r1).strip()
        results['dir4_in_dir1'] = output1
        r2 = env.controller.run_bash_script('test -d /home/user/Desktop/dir3/dir4 && echo YES || echo NO', timeout=30)
        output2 = r2.get('output', r2.get('stdout', '')).strip() if isinstance(r2, dict) else str(r2).strip()
        results['dir4_in_dir3'] = output2
    except Exception as e:
        return {'error': str(e)}
    return results

def get_import_file__c191ddd519b4043306a4effe99ae0aac(env, config: dict):
    """Read the import statements file from VM and return parsed content."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = [l.strip() for l in content.strip().split('\n') if l.strip()]
    import_lines = [l for l in lines if l.startswith('import ') or l.startswith('from ')]
    non_import_lines = [l for l in lines if not (l.startswith('import ') or l.startswith('from '))]
    return {'content': content, 'lines': lines, 'line_count': len(lines), 'import_lines': import_lines, 'non_import_lines': non_import_lines, 'is_sorted': lines == sorted(lines), 'has_duplicates': len(set(lines)) != len(lines)}

def get_dir_structure__f789180956e4b0547b5929aa18a0352e(env, config: dict):
    """Check existence of multiple directories under a base path."""
    base_path = config.get('base_path', '/home/user/project')
    subdirs = config.get('subdirs', [])
    results = {}
    for subdir in subdirs:
        full_path = f'{base_path}/{subdir}'
        cmd_result = env.controller.run_bash_script(f'test -d "{full_path}" && echo "exists" || echo "missing"', timeout=10)
        output = cmd_result.get('output', '').strip() if isinstance(cmd_result, dict) else str(cmd_result).strip()
        results[subdir] = output == 'exists'
    return results

def get_txt_file_content__f7a75c940bb0ba7687c094e66be96f16(env, config: dict):
    """Read a text file from the VM and return its full content and lines."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': '', 'lines': []}
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return {'content': content.strip(), 'lines': lines, 'count': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_text_file_content__55e51bced75ba3e5e08e1be190df564e(env, config: dict):
    """Read a text file from the VM and return its stripped content."""
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File not found: {path}')
            return {'error': 'File not found', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return {'error': str(e), 'content': ''}

def get_text_file_content__358385204fa7b84b2d9118b3314be461(env, config: dict):
    """Read a text file from VM and return its content."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        return {'error': str(e), 'content': ''}

def get_txt_file_lines__5e14eceba731d7ebe21b244525a2edc1(env, config: dict):
    """Read a text file from the VM and return its non-empty lines."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'lines': []}
    with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return {'lines': lines, 'count': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_dir_ownership__425163a454302c61dc411d05ac9420b7(env, config: dict):
    """Get directory ownership and permissions."""
    dirpath = config.get('path', '/home/test1')
    try:
        result = env.controller.run_bash_script(f"stat -c '%U %G %a' {dirpath} 2>/dev/null", timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        if not output:
            return {'error': 'Directory not found or stat failed'}
        parts = output.split()
        if len(parts) >= 3:
            return {'owner': parts[0], 'group': parts[1], 'permissions': parts[2]}
        return {'error': f'Unexpected stat output: {output}'}
    except Exception as e:
        logger.error(f'Error getting directory info: {e}')
        return {'error': str(e)}

def get_text_file_content__ae954eaf7f0ec05c9c2918c4eec73bcc(env, config: dict):
    """Read a text file from the VM and return its stripped content."""
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File not found: {path}')
            return {'error': 'File not found', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace').strip()
        return {'content': content}
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return {'error': str(e), 'content': ''}

def get_total_file_content__aa615bdde918e746d9793d851a2bb4f8(env, config: dict):
    """Read the content of total.txt from the Desktop."""
    file_path = config.get('path', '/home/user/Desktop/total.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    return {'content': content}

def get_folder_file_list__1a224fb0890014daa398af667e7e35c0(env, config: dict):
    """List files in the specified folder on the VM."""
    folder_path = config.get('path', '')
    result = env.controller.run_bash_script(f'ls "{folder_path}" 2>/dev/null', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if not output:
        return {'files': [], 'raw': ''}
    files = sorted([f for f in output.split('\n') if f.strip()])
    return {'files': files, 'raw': output}

def get_file_system_check__f073ea69e2970c4ed8abb5029189b0ac(env, config: dict):
    """Check if specified directory and file exist on the VM."""
    try:
        dir_path = config.get('dir_path', '')
        file_path = config.get('file_path', '')
        dir_result = env.controller.run_bash_script(f'test -d "{dir_path}" && echo "yes" || echo "no"', timeout=30)
        dir_output = dir_result.get('output', dir_result.get('stdout', '')).strip()
        file_result = env.controller.run_bash_script(f'test -f "{file_path}" && echo "yes" || echo "no"', timeout=30)
        file_output = file_result.get('output', file_result.get('stdout', '')).strip()
        return {'dir_exists': dir_output == 'yes', 'file_exists': file_output == 'yes'}
    except Exception as e:
        return {'error': str(e)}

def get_vm_file_content__28ae802c5923a6e7ecfebd0e5cc59702(env, config: dict):
    """Get file content from VM for mpstat per-core CPU report."""
    file_path = config.get('path', '')
    dest = config.get('dest', 'result.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        cache_dir = env.cache_dir if hasattr(env, 'cache_dir') else '/tmp'
        dest_path = os.path.join(cache_dir, dest)
        with open(dest_path, 'wb') as f:
            f.write(file_bytes)
        with open(dest_path, 'r') as f:
            content = f.read()
        return {'content': content, 'file_path': dest_path}
    except Exception as e:
        logger.error(f'get_vm_file_content error: {e}')
        return {'error': str(e)}

def get_sorted_invoice_folders__5d7cb95eaab32c834b0336c7360277e1(env, config: dict):
    """List PDF files in both 'matching' and 'discrepant' folders on Desktop."""
    matching_path = config.get('matching_path', '/home/user/Desktop/matching')
    discrepant_path = config.get('discrepant_path', '/home/user/Desktop/discrepant')
    result_matching = env.controller.run_bash_script(f'ls -1 "{matching_path}" 2>/dev/null', timeout=30)
    stdout_matching = result_matching.get('output', '') if isinstance(result_matching, dict) else str(result_matching)
    matching_files = [f.strip() for f in stdout_matching.strip().split('\n') if f.strip()]
    result_discrepant = env.controller.run_bash_script(f'ls -1 "{discrepant_path}" 2>/dev/null', timeout=30)
    stdout_discrepant = result_discrepant.get('output', '') if isinstance(result_discrepant, dict) else str(result_discrepant)
    discrepant_files = [f.strip() for f in stdout_discrepant.strip().split('\n') if f.strip()]
    return {'matching_files': matching_files, 'discrepant_files': discrepant_files}

def get_dir_structure__95c704f58705a26de50af292b2eb5059(env, config: dict):
    """Get directory structure info from VM."""
    base_path = config.get('path', '')
    subdirs = config.get('subdirs', [])
    try:
        result = env.controller.run_bash_script(f'test -d "{base_path}" && echo "BASE_EXISTS" || echo "BASE_MISSING"', timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        base_exists = 'BASE_EXISTS' in output
        subdir_status = {}
        for sd in subdirs:
            full_path = f'{base_path}/{sd}'
            res = env.controller.run_bash_script(f'test -d "{full_path}" && echo "EXISTS" || echo "MISSING"', timeout=30)
            out = res.get('output', '').strip() if isinstance(res, dict) else str(res).strip()
            subdir_status[sd] = 'EXISTS' in out
        return {'base_exists': base_exists, 'subdirs': subdir_status}
    except Exception as e:
        return {'error': str(e), 'base_exists': False, 'subdirs': {}}

def get_manifest_json__61a5e99298254bc3ea7017f23352ee68(env, config: dict):
    """Read and parse manifest.json from the VM."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {'manifest': data}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_file_info__fda9115554280abf401514b41e43c7b7(env, config: dict):
    """Check if a file exists on the VM and return its info."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'size': 0}
        return {'exists': True, 'size': len(file_bytes), 'first_bytes': file_bytes[:8].hex()}
    except Exception:
        return {'exists': False, 'size': 0}

def get_frame_extract_files__5e444f1b525a6375eead73fcf3fce163(env, config: dict):
    """Check if video frames were extracted to the specified directory as PNG files."""
    extract_dir = config.get('extract_dir', '/tmp/frame_extract')
    result = env.controller.run_bash_script(f'ls {extract_dir}/*.png 2>/dev/null | wc -l', timeout=30)
    stdout = result.get('output', '0').strip() if isinstance(result, dict) else '0'
    try:
        count = int(stdout)
    except (ValueError, TypeError):
        count = 0
    names_result = env.controller.run_bash_script(f'ls {extract_dir}/*.png 2>/dev/null | head -1 && ls {extract_dir}/*.png 2>/dev/null | tail -1', timeout=30)
    names_out = names_result.get('output', '').strip() if isinstance(names_result, dict) else ''
    lines = names_out.split('\n')
    first_frame = lines[0].strip() if len(lines) > 0 else ''
    last_frame = lines[-1].strip() if len(lines) > 1 else first_frame
    return {'frame_count': count, 'first_frame': first_frame, 'last_frame': last_frame, 'extract_dir': extract_dir}

def get_pdf_file_check__1c15db0bd188a0b79dfb455b9076c68e(env, config: dict):
    """Check if a PDF file exists on the VM and has valid content."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'size': 0}
    is_pdf = file_bytes[:4] == b'%PDF'
    return {'exists': True, 'size': len(file_bytes), 'is_pdf': is_pdf}

def get_vm_file_content__a8cfa9828c8ad4ee330aec94e4adaf4c(env, config: dict):
    """Get file content from VM for sar memory report."""
    file_path = config.get('path', '')
    dest = config.get('dest', 'result.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        cache_dir = env.cache_dir if hasattr(env, 'cache_dir') else '/tmp'
        dest_path = os.path.join(cache_dir, dest)
        with open(dest_path, 'wb') as f:
            f.write(file_bytes)
        with open(dest_path, 'r') as f:
            content = f.read()
        return {'content': content, 'file_path': dest_path}
    except Exception as e:
        logger.error(f'get_vm_file_content error: {e}')
        return {'error': str(e)}

def get_docx_text_content__5cd9f8324e43d00ab6826873ca723208(env, config: dict):
    """Get all text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'text': '', 'paragraphs': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        full_text = '\n'.join(paragraphs)
        return {'text': full_text, 'paragraphs': paragraphs, 'paragraph_count': len(paragraphs)}
    finally:
        os.unlink(tmp_path)

def get_lecture_pdf_exists__29ea2d20dd9405b64f3df8f5049b79d3(env, config: dict):
    """Check if a specific lecture PDF exists in the lecture_slides folder."""
    file_path = config.get('path', '/home/user/lecture_slides/lecture1.pdf')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes and len(file_bytes) > 0:
            return {'exists': True, 'size': len(file_bytes)}
        else:
            return {'exists': False, 'size': 0}
    except Exception:
        return {'exists': False, 'size': 0}

def get_stats_file__7636602600807a890f06c684e5901e16(env, config: dict):
    """Read the code stats file from VM and return parsed content."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = [l.strip() for l in content.strip().split('\n') if l.strip()]
    parsed_numbers = []
    for l in lines:
        try:
            parsed_numbers.append(int(l))
        except ValueError:
            parsed_numbers.append(None)
    return {'content': content, 'lines': lines, 'line_count': len(lines), 'parsed_numbers': parsed_numbers}

def get_docx_line_content__31146efdc87a4defa3f97309b7095c90(env, config: dict):
    """Read a docx file and find a line matching a prefix pattern."""
    try:
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/Public Lecture Teaching Plan.docx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs]
            target_prefix = config.get('target_prefix', 'Duration:')
            result = {'paragraphs': paragraphs, 'matched_line': None}
            for text in paragraphs:
                if text.startswith(target_prefix):
                    result['matched_line'] = text
                    break
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pdf_page_info__075106d955ead60b31af19408ecd54ac(env, config: dict):
    """Get PDF page info (dimensions, orientation, margin detection) from a file on the VM."""
    import tempfile, os
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        try:
            from pypdf import PdfReader
        except ImportError:
            return {'error': 'PyPDF2/pypdf not available'}
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        reader = PdfReader(tmp_path)
        num_pages = len(reader.pages)
        if num_pages == 0:
            return {'num_pages': 0, 'width': 0, 'height': 0, 'is_landscape': False, 'file_exists': True, 'has_no_margins': False, 'margin_check_available': False}
        page = reader.pages[0]
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        text_positions = []
        margin_check_available = False
        try:

            def get_text_position__075106d955ead60b31af19408ecd54ac(text, cm, tm, font_dict, font_size):
                if text and text.strip():
                    text_positions.append((float(tm[4]), float(tm[5])))
            page.extract_text(visitor_text=get_text_position__075106d955ead60b31af19408ecd54ac)
            margin_check_available = len(text_positions) > 0
        except Exception:
            pass
        has_no_margins = False
        if margin_check_available:
            min_x = min((p[0] for p in text_positions))
            max_x = max((p[0] for p in text_positions))
            margin_threshold = 30.0
            left_margin = min_x
            right_margin = width - max_x
            has_no_margins = left_margin < margin_threshold and right_margin < margin_threshold
        return {'num_pages': num_pages, 'width': width, 'height': height, 'is_landscape': width > height, 'file_exists': True, 'has_no_margins': has_no_margins, 'margin_check_available': margin_check_available}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_file_content__a4d1b434aaedff8ce21d5e5d82eac67e(env, config: dict):
    """Read text file content from VM and split into lines."""
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': '', 'lines': [], 'line_count': 0}
    try:
        content = file_bytes.decode('utf-8').strip()
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        return {'content': content, 'lines': lines, 'line_count': len(lines)}
    except Exception as e:
        return {'error': str(e), 'content': '', 'lines': [], 'line_count': 0}

def get_file_content__6790a81b7f7c76a950ade3c3fcc6cc45(env, config: dict):
    """Read content of a text file on the VM."""
    file_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'test -f "{file_path}" && echo "exists" || echo "missing"', timeout=30)
    exists_output = exists_result.get('output', '').strip() if isinstance(exists_result, dict) else ''
    if 'missing' in exists_output or not exists_output:
        return {'content': '', 'exists': False}
    result = env.controller.run_bash_script(f'cat "{file_path}" 2>/dev/null', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    lines = [l.strip() for l in output.split('\n') if l.strip()]
    return {'content': output, 'lines': sorted(lines), 'count': len(lines), 'exists': True}

def get_manifest_json__783369d0e7e20c29051ab5abdc75ebb3(env, config: dict):
    """Read and parse manifest.json from the VM."""
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {'manifest': data}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_pdf_on_desktop__18b0663c1e76805fd4ea486cbf72380b(env, config: dict):
    """Check if a PDF file exists at the specified path on the VM."""
    import os
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes and len(file_bytes) > 0:
            return {'exists': True, 'size': len(file_bytes), 'path': path}
        else:
            return {'exists': False, 'size': 0, 'path': path}
    except Exception as e:
        return {'exists': False, 'size': 0, 'path': path, 'error': str(e)}

def get_file_content_check__2372f00ad964142d3685a78ced2bb21c(env, config: dict):
    """Read a file from the VM and return its content for verification."""
    file_path = config.get('file_path', '/home/user/system_info.txt')
    cmd_result = env.controller.run_bash_script(f'cat "{file_path}" 2>/dev/null || echo "__FILE_NOT_FOUND__"', timeout=10)
    output = cmd_result.get('output', '').strip() if isinstance(cmd_result, dict) else str(cmd_result).strip()
    if output == '__FILE_NOT_FOUND__':
        return {'error': 'File not found', 'content': ''}
    return {'content': output}

def get_pdf_file_info__beb4b25746cd97a98dc95327476aaf4d(env, config: dict):
    """Get PDF file info from VM to verify export."""
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/video_tips.pdf')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        is_pdf = file_bytes[:5] == b'%PDF-'
        page_count = 0
        if is_pdf:
            try:
                import subprocess
                result = subprocess.run(['python3', '-c', f"from PyPDF2 import PdfReader; r = PdfReader('{tmp_path}'); print(len(r.pages))"], capture_output=True, text=True, timeout=10)
                if result.returncode == 0:
                    page_count = int(result.stdout.strip())
            except Exception:
                content = file_bytes.decode('latin-1', errors='ignore')
                page_count = content.count('/Type /Page') - content.count('/Type /Pages')
                if page_count <= 0:
                    page_count = content.count('/Type/Page') - content.count('/Type/Pages')
        return {'exists': True, 'is_pdf': is_pdf, 'size': len(file_bytes), 'page_count': page_count}
    finally:
        os.unlink(tmp_path)

def get_vm_files_multi__8aa2e8e2be1ca5d610e977b9c78060b2(env, config: dict):
    """Read multiple files from VM and return their contents as a dict.

    Config:
        files: list of dicts with 'path' and 'key' fields
    Returns:
        dict mapping key -> file content string (or None if not found)
    """
    result = {}
    for file_info in config.get('files', []):
        try:
            file_bytes = env.controller.get_file(file_info['path'])
            if file_bytes:
                result[file_info['key']] = file_bytes.decode('utf-8', errors='replace')
            else:
                result[file_info['key']] = None
        except Exception:
            result[file_info['key']] = None
    return result

def get_result_file_content__4c5e4ba8ba69ff54d1e2443439646d44(env, config: dict):
    """Read a text file from VM and return its content."""
    file_path = config.get('path', '/home/user/Desktop/result.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8').strip()
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1').strip()
    return {'content': content}

def get_csv_head__8b9d6ae5a27e52886aef51e39cc4df8a(env, config: dict):
    """Read first N lines of a CSV file from VM."""
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/file1.csv')
    n_lines = config.get('n_lines', 6)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False, mode='wb') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            lines = []
            for (i, line) in enumerate(f):
                if i >= n_lines:
                    break
                lines.append(line.strip())
        return {'lines': lines, 'line_count': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_csv_content__f43838d5904be041721c003e00fd36b5(env, config: dict):
    """Read first N lines of a CSV file from VM."""
    import tempfile
    import os
    path = config.get('path', '')
    n_lines = config.get('n_lines', 6)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.csv', delete=False, mode='wb') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='replace') as f:
            lines = []
            for (i, line) in enumerate(f):
                if i >= n_lines:
                    break
                lines.append(line.strip())
        return {'lines': lines, 'line_count': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_check_output_file__f4c7cbbe6d6bd73cc0b521d770fb3945(env, config: dict):
    """Get the content of output.txt from VM."""
    file_path = config.get('path', '/home/user/Desktop/output.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    lines = [l.strip() for l in content.splitlines() if l.strip()]
    return {'exists': True, 'content': content, 'lines': lines, 'line_count': len(lines)}

def get_file_info__5d2d411fc8203609cb2ed6eef9424f86(env, config: dict):
    """Download file from VM and check its existence and format."""
    from PIL import Image
    path = config.get('path', '/home/user/Desktop/character.jpg')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'path': path}
    with tempfile.NamedTemporaryFile(suffix=os.path.splitext(path)[1], delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'exists': True, 'format': img.format, 'width': img.width, 'height': img.height, 'file_size': len(file_bytes)}
    except Exception as e:
        return {'exists': True, 'format': 'unknown', 'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_result_file_content__4128b5119649f756a2c7d9fc4321d25f(env, config: dict):
    """Read a text file from VM and return its content."""
    file_path = config.get('path', '/home/user/Desktop/result.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8').strip()
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1').strip()
    return {'content': content}

def get_merged_novel_file__2abe2760620ddfe22faad3f384e1cae4(env, config: dict):
    """Get the merged novel text file from VM and return its line count and content markers."""
    file_path = config.get('path', '/home/user/Documents/Novels/Pass Through/full_novel.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'exists': False}
        content = file_bytes.decode('utf-8', errors='replace')
        lines = content.strip().split('\n')
        has_ch1 = 'Chapter 1' in content or 'Passing through the book' in content
        has_ch2 = 'Chapter 2' in content or 'Change' in content
        has_ch3 = 'Chapter 3' in content or 'Plan' in content
        has_ch4 = 'Chapter 4' in content or 'Black Material' in content
        has_ch5 = 'Chapter 5' in content or 'Ask him to pay back the money' in content
        return {'exists': True, 'line_count': len(lines), 'byte_count': len(file_bytes), 'has_ch1': has_ch1, 'has_ch2': has_ch2, 'has_ch3': has_ch3, 'has_ch4': has_ch4, 'has_ch5': has_ch5}
    except Exception as e:
        logger.error(f'Error reading merged novel file: {e}')
        return {'error': str(e), 'exists': False}

def get_local_folder_list__91f316cac3c13e1ffa3f1d0a1879bf2b(env, config: dict):
    """List files in Thunderbird Local Folders directory and save to a temp file."""
    profile_dir = config.get('profile_dir', '/home/user/.thunderbird/t5q2a5hp.default-release/')
    local_folders_dir = os.path.join(profile_dir, 'Mail', 'Local Folders')
    cmd = f"ls -1 '{local_folders_dir}' 2>/dev/null"
    result = env.controller.run_bash_script(cmd, timeout=10)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '')
    elif isinstance(result, str):
        output = result
    tmp_path = os.path.join(tempfile.gettempdir(), 'local_folders_list.txt')
    with open(tmp_path, 'w') as f:
        f.write(output)
    return tmp_path

def get_desktop_file__bc88c1a76af0c575a8c22634fe96ef8b_qw35sft2_f69c4fb7(env, config: dict):
    """Check if a specific file exists on the Desktop."""
    filename = config.get('filename', '')
    if not filename:
        return {'exists': False, 'filename': '', 'error': 'no filename in config'}
    path = f'/home/user/Desktop/{filename}'
    res = env.controller.run_bash_script(f'test -f "{path}" && echo "exists" || echo "not_found"', timeout=10)
    output = res.get('output', '').strip() if isinstance(res, dict) else ''
    return {'exists': output == 'exists', 'filename': filename, 'path': path}

def get_downloads_pdf__5416755d2dd223fd56c07a64ea26e507_qw35sft2_5814a28c(env, config: dict):
    """Get PDF files in Downloads folder and page count for no-margin verification."""
    res = env.controller.run_bash_script('ls /home/user/Downloads/ 2>/dev/null', timeout=10)
    output = res.get('output', '') if isinstance(res, dict) else ''
    files = [f.strip() for f in output.strip().split('\n') if f.strip()]
    pdf_files = [f for f in files if f.lower().endswith('.pdf')]
    result = {'pdf_files': pdf_files, 'count': len(pdf_files)}
    name_filter = config.get('pdf_name_contains', 'llm').lower()
    target_pdf = next((f for f in pdf_files if name_filter in f.lower()), None)
    if target_pdf:
        try:
            file_bytes = env.controller.get_file(f'/home/user/Downloads/{target_pdf}')
            if file_bytes:
                tmp_path = None
                try:
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    import fitz
                    doc = fitz.open(tmp_path)
                    result['page_count'] = len(doc)
                    doc.close()
                finally:
                    if tmp_path and os.path.exists(tmp_path):
                        os.unlink(tmp_path)
        except Exception as e:
            result['page_count_error'] = str(e)
    return result

def get_file_exists__de94d61f074111bc63d9a066e03aa46f_qw35sft2_73b15751(env, config: dict):
    """Check whether clip.gif is a valid animated GIF with ~2 seconds of content."""
    path = config.get('path', '/home/user/clip.gif')
    result = {'exists': False, 'is_gif': False, 'frame_count': 0, 'duration_secs': 0.0}
    r = env.controller.run_bash_script(f'test -f "{path}" && echo FILE_EXISTS_OK || echo FILE_MISSING', timeout=10)
    if not r or 'FILE_EXISTS_OK' not in r.get('output', ''):
        return result
    result['exists'] = True
    r = env.controller.run_bash_script(f'file "{path}" 2>/dev/null', timeout=10)
    file_output = r.get('output', '') if r else ''
    result['is_gif'] = 'GIF' in file_output
    if not result['is_gif']:
        return result
    r = env.controller.run_bash_script(f'identify "{path}" 2>/dev/null | wc -l', timeout=30)
    try:
        result['frame_count'] = int((r.get('output', '') if r else '').strip())
    except (ValueError, AttributeError):
        result['frame_count'] = 0
    awk_sum = "awk '{s+=$1} END {print s+0}'"
    r = env.controller.run_bash_script(f'identify -format "%T\\n" "{path}" 2>/dev/null | {awk_sum}', timeout=30)
    try:
        total_cs = float((r.get('output', '') if r else '0').strip())
        result['duration_secs'] = total_cs / 100.0
    except (ValueError, AttributeError):
        result['duration_secs'] = 0.0
    if result['duration_secs'] == 0.0 and result['frame_count'] > 0:
        r = env.controller.run_bash_script(f'ffprobe -v error -select_streams v:0 -show_entries stream=duration -of csv=p=0 "{path}" 2>/dev/null', timeout=30)
        try:
            result['duration_secs'] = float((r.get('output', '') if r else '0').strip())
        except (ValueError, AttributeError):
            result['duration_secs'] = 0.0
    return result

def get_csv_and_pdf__f45249b0314207a95d7366d15bee2907_qw35sft2_1e85c06e(env, config: dict):
    """Check both Export_Calc_to_CSV.csv and Export_Calc_to_CSV.pdf exist on the VM."""
    result = {'csv_exists': False, 'csv_row_count': 0, 'pdf_exists': False, 'pdf_size': 0}
    csv_bytes = env.controller.get_file('/home/user/Export_Calc_to_CSV.csv')
    if csv_bytes:
        result['csv_exists'] = True
        content = csv_bytes.decode('utf-8', errors='replace')
        lines = [l.strip() for l in content.strip().splitlines() if l.strip()]
        result['csv_row_count'] = len(lines)
    pdf_bytes = env.controller.get_file('/home/user/Export_Calc_to_CSV.pdf')
    if pdf_bytes:
        result['pdf_exists'] = True
        result['pdf_size'] = len(pdf_bytes)
    return result

def get_bold_header_and_pdf__86ea3ac794566dae35e2163b9c7d3c94_qw35sft2_e37cd2f8(env, config: dict):
    """Check bold formatting of entire header row 11 (A11, B11, C11), fit-to-page setting, and PDF existence."""
    pdf_bytes = env.controller.get_file('/home/user/Resize_Cells_Fit_Page.pdf')
    pdf_exists = bool(pdf_bytes and len(pdf_bytes) > 0)
    xlsx_path = config.get('path', '/home/user/Resize_Cells_Fit_Page.xlsx')
    file_bytes = env.controller.get_file(xlsx_path)
    if not file_bytes:
        return {'error': 'xlsx not found', 'pdf_exists': pdf_exists, 'header_bold': False, 'fit_to_page': False}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        a11_bold = bool(ws['A11'].font and ws['A11'].font.bold)
        b11_bold = bool(ws['B11'].font and ws['B11'].font.bold)
        c11_bold = bool(ws['C11'].font and ws['C11'].font.bold)
        header_bold = a11_bold and b11_bold and c11_bold
        try:
            page_setup_pr = ws.sheet_properties.pageSetUpPr
            fit_to_page = bool(page_setup_pr is not None and page_setup_pr.fitToPage)
        except Exception:
            fit_to_page = False
        return {'pdf_exists': pdf_exists, 'header_bold': header_bold, 'fit_to_page': fit_to_page}
    finally:
        os.unlink(tmp_path)

def get_col_a_bold_and_csv__8f6a3ccc5fb938975de38e8fbcc8f580_qw35sft2_725bcefe(env, config: dict):
    """Read xlsx column A bold state and check CSV existence for bold+export task."""
    import tempfile, os
    try:
        import openpyxl
    except ImportError:
        return {'error': 'openpyxl not available', 'all_col_a_bold': False, 'csv_exists': False}
    result = {'all_col_a_bold': False, 'col_a_bold': [], 'csv_exists': False, 'csv_row_count': 0}
    csv_bytes = env.controller.get_file('/home/user/Export_Calc_to_CSV.csv')
    if csv_bytes:
        result['csv_exists'] = True
        content = csv_bytes.decode('utf-8', errors='replace')
        lines = [l.strip() for l in content.strip().splitlines() if l.strip()]
        result['csv_row_count'] = len(lines)
    xlsx_bytes = env.controller.get_file('/home/user/Export_Calc_to_CSV.xlsx')
    if xlsx_bytes:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(xlsx_bytes)
                tmp_path = tmp.name
            wb = openpyxl.load_workbook(tmp_path)
            ws = wb.worksheets[0]
            bold_flags = []
            for row_idx in range(1, 5):
                cell = ws.cell(row=row_idx, column=1)
                bold_flags.append(bool(cell.font and cell.font.bold))
            result['col_a_bold'] = bold_flags
            result['all_col_a_bold'] = all(bold_flags)
        except Exception:
            result['col_a_bold'] = []
            result['all_col_a_bold'] = False
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)
    return result

def get_file_exists_on_desktop__3ed023e79bfe94e32e5ea286856ad04c_qw35sft2_b66cefae(env, config: dict):
    """Check if an exported file exists at the given path on the VM."""
    filepath = config.get('path', '')
    result = env.controller.run_bash_script(f'test -f "{filepath}" && echo "exists" || echo "not_found"', timeout=15)
    if isinstance(result, dict):
        output = result.get('output', '') or result.get('stdout', '') or ''
    else:
        output = str(result) if result else ''
    return {'exists': 'exists' in output.strip(), 'path': filepath}

def get_pdf_exists__19cbc30a1547517beef14a189cce767e_qw35sft2_652518fc(env, config: dict):
    """Check if View_Person_Organizational_Summary.pdf exists on the Desktop."""
    path = config.get('path', '/home/user/Desktop/View_Person_Organizational_Summary.pdf')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'size': 0}
    return {'exists': True, 'size': len(file_bytes)}

def get_pdf_and_odt_exists__c65906d6f690144adc013e7fcea2301e_qw35sft2_5e4aefbc(env, config: dict):
    """Check if both the PDF and ODT copies exist on the Desktop."""
    pdf_path = config.get('pdf_path', '/home/user/Desktop/View_Person_Organizational_Summary.pdf')
    odt_path = config.get('odt_path', '/home/user/Desktop/View_Person_Organizational_Summary.odt')
    pdf_bytes = env.controller.get_file(pdf_path)
    odt_bytes = env.controller.get_file(odt_path)
    return {'pdf_exists': bool(pdf_bytes and len(pdf_bytes) > 0), 'odt_exists': bool(odt_bytes and len(odt_bytes) > 0)}

def get_para0_word_content__469a4ad5eae55be7063b8ec35b77b37d_qw35sft2_947cb49d(env, config: dict):
    """Get text content of the first paragraph to verify word replacement."""
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/CCCH9003_Tutorial_guidelines.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        text = para.text
        return {'text': text, 'has_talkative': 'talkative' in text.lower(), 'has_participative': 'participative' in text.lower()}
    finally:
        os.unlink(tmp_path)

def get_orgsummary_pdf_exists__82b855263e1bdeb6cb9c2f640e53bf6c_qw35sft2_1748291e(env, config: dict):
    """Check if OrgSummary.pdf exists on the Desktop."""
    path = config.get('path', '/home/user/Desktop/OrgSummary.pdf')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'size': 0}
    return {'exists': True, 'size': len(file_bytes)}

def get_pdf_in_documents__df1f48650ab6f100a1208b8040cd8828_qw35sft2_6e48e714(env, config: dict):
    """Check if View_Person_Organizational_Summary.pdf exists in the Documents folder."""
    path = config.get('path', '/home/user/Documents/View_Person_Organizational_Summary.pdf')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'size': 0}
    return {'exists': True, 'size': len(file_bytes)}

def get_pdf_and_footer_pagenum__0bb862e1956085267ced241046059b10_qw35sft2_f87e83a2(env, config: dict):
    """Check PDF existence and whether the docx footer contains page number fields."""
    pdf_path = config.get('pdf_path', '/home/user/Desktop/View_Person_Organizational_Summary.pdf')
    docx_path = config.get('docx_path', '/home/user/Desktop/View_Person_Organizational_Summary.docx')
    pdf_bytes = env.controller.get_file(pdf_path)
    pdf_exists = bool(pdf_bytes and len(pdf_bytes) > 0)
    has_page_numbers = False
    docx_bytes = env.controller.get_file(docx_path)
    if docx_bytes:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(docx_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            for section in doc.sections:
                footer = section.footer
                xml_str = footer._element.xml
                if 'PAGE' in xml_str and 'fldChar' in xml_str or 'instrText' in xml_str:
                    has_page_numbers = True
                    break
        except Exception:
            pass
        finally:
            os.unlink(tmp_path)
    return {'pdf_exists': pdf_exists, 'has_page_numbers': has_page_numbers}

def get_path_text_file__a43735a3f0450074dfa18129f92a73a9_qw35sft2_74219b39(env, config: dict):
    """Read a text file's contents from the VM filesystem."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found or empty', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    return {'content': content}

def get_text_file_content__f36e167c6b956664e01e5e47712b671b_qw35sft2_c7971e50(env, config: dict):
    """Read a text file from the VM and return its stripped content."""
    path = config.get('path', '/home/user/chapter_count.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        return {'content': file_bytes.decode('utf-8').strip()}
    except Exception as e:
        return {'error': str(e), 'content': ''}

def get_file_content__5ca68902218b7e30665245f17696931d_qw35sft2_6875fbcc(env, config: dict):
    """Read the contents of a file on the VM and return its lines."""
    file_path = config.get('path', '/home/user/Test/Speed/results.txt')
    script = f'cat {file_path!r} 2>/dev/null || echo "__FILE_NOT_FOUND__"'
    try:
        result = env.controller.run_bash_script(script, timeout=15)
        if isinstance(result, dict):
            output = result.get('output', result.get('stdout', ''))
        else:
            output = str(result) if result else ''
        if '__FILE_NOT_FOUND__' in output:
            return {'error': 'File not found', 'lines': [], 'content': ''}
        lines = [ln.rstrip('\r') for ln in output.split('\n')]
        non_empty = [ln for ln in lines if ln.strip()]
        return {'lines': non_empty, 'content': output, 'line_count': len(non_empty)}
    except Exception as e:
        return {'error': str(e), 'lines': [], 'content': ''}

def get_desktop_file_exists__6964660b6570aa960a367d963aeb2477_qw35sft2_50044e66(env, config: dict):
    """Check if a specific file exists on the Desktop and whether original location is clear."""
    dest_path = config.get('path', '/home/user/Desktop/1. Symmetric matrices and adjacency of a graph.pdf')
    orig_path = config.get('orig_path', '/home/user/Desktop/book/1. Symmetric matrices and adjacency of a graph.pdf')
    result_dest = env.controller.run_bash_script(f'test -f "{dest_path}" && echo "exists" || echo "not_found"', timeout=15)
    out_dest = result_dest.get('output', '') if isinstance(result_dest, dict) else str(result_dest)
    result_orig = env.controller.run_bash_script(f'test -f "{orig_path}" && echo "exists" || echo "not_found"', timeout=15)
    out_orig = result_orig.get('output', '') if isinstance(result_orig, dict) else str(result_orig)
    return {'file_at_dest': 'exists' in out_dest, 'file_at_orig': 'exists' in out_orig}

def get_dir_exists__e66cc6b0ec6479dc303b558a9cc72642_qw35sft2_e41816cf(env, config: dict):
    """Check if a directory exists on the VM."""
    dir_path = config.get('dir_path', '')
    result = env.controller.run_bash_script(f'[ -d "{dir_path}" ] && echo "exists" || echo "not_found"', timeout=10)
    if isinstance(result, dict):
        output = result.get('output', '').strip()
    else:
        output = str(result).strip()
    return {'exists': output == 'exists', 'path': dir_path}

def get_res_txt_content__da27b9fa71ba67953d0b9f649a3197fa_qw35sft2_ead3d6e3(env, config: dict):
    """Read res.txt from the VM Desktop and return its content."""
    file_bytes = env.controller.get_file('/home/user/Desktop/res.txt')
    if not file_bytes:
        return {'error': 'File not found: /home/user/Desktop/res.txt'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {e}'}
    return {'content': content}

def get_receipts_file_listing__d248703d35223088b5508b0c05332ea3_qw35sft2_cbb3b1f6(env, config: dict):
    """List files in the receipts folder and return as a newline-separated string."""
    result = env.controller.run_bash_script('ls /home/user/Documents/Finance/receipts/', timeout=10)
    if isinstance(result, dict):
        return result.get('output', '') or result.get('stdout', '')
    return str(result) if result else ''

def get_text_file_content__30cadb0ec7eca28df693cf924ab88301_qw35sft2_a743245c(env, config: dict):
    """Read a text file from the VM and return its content."""
    path = config.get('path', '/home/user/Desktop/image_info.txt')
    result = env.controller.run_bash_script(f'cat {path} 2>/dev/null || echo FILE_NOT_FOUND', timeout=15)
    output = (result.get('output') or '').strip() if isinstance(result, dict) else str(result).strip()
    if output == 'FILE_NOT_FOUND' or not output:
        return {'error': 'File not found or empty', 'content': None}
    return {'content': output}

def get_fs_file_exists__21cbb74de44bf7839b1ea14bdfabea03_qw35sft2_9687eb5c(env, config: dict):
    """Check if a file exists on the VM filesystem."""
    path = config.get('path', '/home/user/Projects/happy-extension/manifest.json')
    result = env.controller.run_bash_script(f'[ -f "{path}" ] && echo "YES" || echo "NO"', timeout=15)
    if isinstance(result, dict):
        output = result.get('output', '') or result.get('stdout', '') or str(result)
    else:
        output = str(result)
    return {'exists': 'YES' in output.strip()}

def get_txt_multihop_gemini__61501dfdaa17fec1d3c116cb7f744d5c_qw35sft2_1fcc099e(env, config: dict):
    """Get content of gemini_multihop.txt from the Desktop."""
    result = env.controller.run_bash_script('cat /home/user/Desktop/gemini_multihop.txt 2>/dev/null || echo "__FILE_NOT_FOUND__"', timeout=10)
    stdout = result.get('stdout', '') if isinstance(result, dict) else ''
    if '__FILE_NOT_FOUND__' in stdout or not stdout.strip():
        return {'error': 'file not found or empty'}
    content = stdout.strip()
    return {'content': content, 'contains_iliad': 'Iliad' in content, 'contains_masculine': 'masculine' in content.lower(), 'char_count': len(content)}

def get_csv_filtered_a__2fadb675c56bf8de284fee43324487b2_qw35sft2_b0731958(env, config: dict):
    """Read filtered CSV (last names starting with 'A') from VM."""
    path = config.get('path', '/home/user/Desktop/filtered.csv')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'path': path}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        return {'error': f'Decode error: {e}'}
    lines = [l.strip() for l in content.splitlines() if l.strip()]
    if not lines:
        return {'error': 'Empty file'}
    first_line = lines[0]
    if 'name' in first_line.lower() or 'first' in first_line.lower() or 'last' in first_line.lower():
        data_lines = lines[1:]
    else:
        data_lines = lines
    all_start_with_a = True
    for line in data_lines:
        if '\t' in line:
            parts = line.split('\t')
            last = parts[1].strip() if len(parts) > 1 else ''
        elif ',' in line:
            row = next(csv.reader(io.StringIO(line)), [])
            last = row[1].strip() if len(row) > 1 else ''
        elif ' ' in line:
            parts = line.rsplit(' ', 1)
            last = parts[-1].strip()
        else:
            last = line.strip()
        if last and (not last.upper().startswith('A')):
            all_start_with_a = False
            break
    return {'data_row_count': len(data_lines), 'total_line_count': len(lines), 'all_last_names_start_with_a': all_start_with_a, 'first_data_line': data_lines[0] if data_lines else ''}

def get_dir_exists__e103d8977cd2fefeec7972ff8169e36d_qw35sft2_0dfba723(env, config: dict):
    """Check whether a directory exists on the VM via bash."""
    dir_path = config.get('path', '/home/user/Test/Speed')
    script = f'test -d {dir_path!r} && echo "exists" || echo "not_exists"'
    try:
        result = env.controller.run_bash_script(script, timeout=15)
        if isinstance(result, dict):
            output = result.get('output', result.get('stdout', ''))
        else:
            output = str(result) if result else ''
        return {'exists': output.strip() == 'exists'}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_file_at_path__82c707062135fc7c567ff519e1b5a575_qw35sft2_273ea9cf(env, config: dict):
    """Check whether a file exists at the target path."""
    target_path = config.get('target_path', '')
    script = f'if [ -f "{target_path}" ]; then echo "exists"; else echo "missing"; fi'
    result = env.controller.run_bash_script(script, timeout=10)
    if result is None:
        return {'error': 'run_bash_script returned None', 'exists': False}
    output = result.get('output', '')
    return {'exists': 'exists' in output}

def get_bubblesort_both_files__8c65eee451a1bb106a548123af80ea08_qw35sft2_60135b55(env, config: dict):
    """Read both bubbleSort.py and res.txt from the VM Desktop."""
    result = {}
    py_bytes = env.controller.get_file('/home/user/Desktop/bubbleSort.py')
    if py_bytes:
        try:
            py_content = py_bytes.decode('utf-8')
            result['py_content'] = py_content
        except Exception:
            result['py_content'] = None
    else:
        result['py_content'] = None
    txt_bytes = env.controller.get_file('/home/user/Desktop/res.txt')
    if txt_bytes:
        try:
            txt_content = txt_bytes.decode('utf-8')
            result['txt_content'] = txt_content
        except Exception:
            result['txt_content'] = None
    else:
        result['txt_content'] = None
    return result

def get_main_py_content__f65942d684aa73fa8a727494a6335877_qw35sft2_9e5513a5(env, config: dict):
    """Read main.py from the VS Code project directory."""
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/project/main.py')
        if isinstance(file_bytes, bytes):
            return file_bytes.decode('utf-8')
        return str(file_bytes)
    except Exception as e:
        return ''

def get_csv_space_merged__00d80996b3209df110a0b1e69fc0ab31_qw35sft2_b3c91bc9(env, config: dict):
    """Read output CSV and check if rows are space-joined full names (single column)."""
    path = config.get('path', '/home/user/Desktop/output.csv')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'path': path}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        return {'error': f'Decode error: {e}'}
    lines = [l.strip() for l in content.splitlines() if l.strip()]
    if len(lines) < 2:
        return {'error': 'Too few lines', 'line_count': len(lines)}
    first_line = lines[0]
    if '\t' in first_line or 'name' in first_line.lower():
        data_lines = lines[1:]
    else:
        data_lines = lines
    if not data_lines:
        return {'error': 'No data lines'}
    first_data = data_lines[0]
    second_data = data_lines[1] if len(data_lines) > 1 else ''
    has_tab = '\t' in first_data
    has_space_format = ' ' in first_data and (not has_tab)
    return {'first_data_line': first_data, 'second_data_line': second_data, 'has_tab': has_tab, 'has_space_format': has_space_format, 'data_row_count': len(data_lines), 'total_line_count': len(lines)}

def get_file_exists__057715c74cd9cb2c93d92377c04bb077_qw35sft2_722a620b(env, config: dict):
    """Check whether a file exists on the VM via bash."""
    file_path = config.get('path', '/home/user/Test/Speed/results.txt')
    script = f'test -f {file_path!r} && echo "exists" || echo "not_exists"'
    try:
        result = env.controller.run_bash_script(script, timeout=15)
        if isinstance(result, dict):
            output = result.get('output', result.get('stdout', ''))
        else:
            output = str(result) if result else ''
        return {'exists': output.strip() == 'exists'}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_clipboard_content__241466d2f7566981ed09035f56716bb4_qw35sft2_9c0b84e6(env, config: dict):
    """Read the current X clipboard content via xsel (with xclip fallback)."""
    script = "xsel --clipboard --output 2>/dev/null || xclip -selection clipboard -o 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(script, timeout=10)
    if result is None:
        return {'error': 'run_bash_script returned None', 'clipboard': ''}
    if result.get('returncode', 0) != 0 and (not result.get('output', '')):
        return {'error': result.get('error', ''), 'clipboard': ''}
    content = result.get('output', '').strip()
    return {'clipboard': content}

def get_res_txt_content__68c8947721f1221211aba42cd6d1735d_qw35sft2_37d291c9(env, config: dict):
    """Read res.txt from the VM Desktop and return its content."""
    file_bytes = env.controller.get_file('/home/user/Desktop/res.txt')
    if not file_bytes:
        return {'error': 'File not found: /home/user/Desktop/res.txt'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {e}'}
    return {'content': content}

def get_file_state__ea8c7a7a44f5cb78e25fd89d658a863f_qw35sft2_fc8fdbc1(env, config: dict):
    """Check if a file exists, has non-zero size, and get its media duration via ffprobe."""
    file_path = config.get('file_path', '')
    result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && stat -c %s "{file_path}" || echo "not_found"', timeout=10)
    if isinstance(result, dict):
        output = result.get('output', '').strip()
    else:
        output = str(result).strip()
    if not output or output == 'not_found':
        return {'exists': False, 'size': 0, 'duration': -1.0}
    try:
        size = int(output)
    except (ValueError, TypeError):
        return {'exists': False, 'size': 0, 'duration': -1.0}
    dur_result = env.controller.run_bash_script(f'ffprobe -v error -show_entries format=duration -of default=noprint_wrappers=1:nokey=1 "{file_path}" 2>/dev/null', timeout=30)
    if isinstance(dur_result, dict):
        dur_output = dur_result.get('output', '').strip()
    else:
        dur_output = str(dur_result).strip()
    try:
        duration = float(dur_output)
    except (ValueError, TypeError):
        duration = -1.0
    return {'exists': True, 'size': size, 'duration': duration}

def get_all_file_permissions__1540d35f6307434ff998913b82a97dfa_qw35sft2_ae6ad5fe(env, config: dict):
    """Get permissions of all regular files under testDir."""
    result = env.controller.run_bash_script('find /home/user/testDir -type f -exec stat -c "%a %n" {} \\; 2>/dev/null', timeout=15)
    lines = [l.strip() for l in result.get('stdout', '').strip().split('\n') if l.strip()]
    permissions = {}
    for line in lines:
        parts = line.split(' ', 1)
        if len(parts) == 2:
            permissions[parts[1]] = parts[0]
    return {'permissions': permissions, 'count': len(permissions)}

def get_user_identity_file__9be4c7481abd67d84bd529b63ed66ca8_qw35sft2_9dcd61e2(env, config: dict):
    """Read ~/user_identity.txt containing output of the `id` command."""
    path = '/home/user/user_identity.txt'
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    return {'content': content}

def get_count_file__5799d2a2e180d6540c452f203d2bde46_qw35sft2_55233ea8(env, config: dict):
    """Check archive in old_files (with correct content), new_files populated, and count.txt."""
    result = env.controller.run_bash_script('ARCH=$(find /tmp/test_files/old_files -maxdepth 1 -name "*.tar.gz" | head -1); COUNT=$(cat /tmp/test_files/count.txt 2>/dev/null | tr -d "[:space:]"); NEW_COUNT=$(ls /tmp/test_files/new_files/ 2>/dev/null | wc -l | tr -d "[:space:]"); if [ -n "$ARCH" ]; then   ARCH_CONTENT=$(tar -tzf "$ARCH" 2>/dev/null | grep -E "old_file[12]\\.txt" | wc -l | tr -d "[:space:]"); else   ARCH_CONTENT=0; fi; echo "ARCHIVE=${ARCH}"; echo "COUNT=${COUNT}"; echo "NEW_COUNT=${NEW_COUNT}"; echo "ARCH_CONTENT=${ARCH_CONTENT}"', timeout=15)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    archive_path = ''
    count_content = ''
    new_count = '0'
    arch_content = '0'
    for line in output.strip().splitlines():
        if line.startswith('ARCHIVE='):
            archive_path = line[len('ARCHIVE='):].strip()
        elif line.startswith('COUNT='):
            count_content = line[len('COUNT='):].strip()
        elif line.startswith('NEW_COUNT='):
            new_count = line[len('NEW_COUNT='):].strip()
        elif line.startswith('ARCH_CONTENT='):
            arch_content = line[len('ARCH_CONTENT='):].strip()
    return {'archive_exists': bool(archive_path), 'archive_has_old_files': int(arch_content) >= 2, 'new_files_populated': int(new_count) > 0, 'count_content': count_content}

def get_copy_4dirs__4f3015d5aa890d8fea505e363d3f7aff_qw35sft2_2fd46445(env, config: dict):
    """Check file1 in dir1, dir2, dir3, and new dir4."""
    result = env.controller.run_bash_script('echo dir1:$(test -f /home/user/dir1/file1 && echo yes || echo no)\necho dir2:$(test -f /home/user/dir2/file1 && echo yes || echo no)\necho dir3:$(test -f /home/user/dir3/file1 && echo yes || echo no)\necho dir4:$(test -f /home/user/dir4/file1 && echo yes || echo no)', timeout=10)
    if not result:
        return {'error': 'command failed'}
    output = result.get('output', '')
    return {'dir1_has_file1': 'dir1:yes' in output, 'dir2_has_file1': 'dir2:yes' in output, 'dir3_has_file1': 'dir3:yes' in output, 'dir4_has_file1': 'dir4:yes' in output}

def get_file_perms_and_rename__653ac9883bb8d908f613f5a6b4d6a405_qw35sft2_68e3d2be(env, config: dict):
    """Get permissions of all regular files and check if file3.txt was renamed to file3_backup.txt."""
    perm_result = env.controller.run_bash_script('find /home/user/testDir -type f -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    backup_result = env.controller.run_bash_script('test -f /home/user/testDir/subDir2/file3_backup.txt && echo "exists" || echo "missing"', timeout=10)
    old_result = env.controller.run_bash_script('test -f /home/user/testDir/subDir2/file3.txt && echo "exists" || echo "missing"', timeout=10)
    file_perms = [p.strip() for p in perm_result.get('stdout', '').strip().split('\n') if p.strip()]
    backup_exists = backup_result.get('stdout', '').strip() == 'exists'
    old_gone = old_result.get('stdout', '').strip() == 'missing'
    return {'file_permissions': file_perms, 'backup_file_exists': backup_exists, 'old_file_gone': old_gone}

def get_sys_info_file__5cefed7a5e164dc00f5c88ae5726c23e_qw35sft2_5a973c8c(env, config: dict):
    """Read ~/sys_info.txt containing whoami and hostname outputs."""
    path = '/home/user/sys_info.txt'
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    return {'content': content}

def get_old_files_cleanup__aba3d770f03c5889eb920c2ccea4f4e0_qw35sft2_1335d4ce(env, config: dict):
    """Check that old_files has an archive but no loose txt files, and that new_files has recent txt files."""
    result = env.controller.run_bash_script('ARCH=$(find /tmp/test_files/old_files -maxdepth 1 -name "*.tar.gz" | head -1); LOOSE=$(find /tmp/test_files/old_files -maxdepth 1 -type f -name "*.txt" | wc -l); NEW_TXT=$(find /tmp/test_files/new_files -maxdepth 1 -type f -name "*.txt" | wc -l); echo "ARCHIVE=${ARCH}"; echo "LOOSE_TXT=${LOOSE}"; echo "NEW_TXT=${NEW_TXT}"', timeout=15)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    archive_path = ''
    loose_count = 0
    new_txt_count = 0
    for line in output.strip().splitlines():
        if line.startswith('ARCHIVE='):
            archive_path = line[len('ARCHIVE='):].strip()
        elif line.startswith('LOOSE_TXT='):
            try:
                loose_count = int(line[len('LOOSE_TXT='):].strip())
            except ValueError:
                pass
        elif line.startswith('NEW_TXT='):
            try:
                new_txt_count = int(line[len('NEW_TXT='):].strip())
            except ValueError:
                pass
    return {'archive_exists': bool(archive_path), 'archive_path': archive_path, 'loose_txt_count': loose_count, 'new_txt_count': new_txt_count}

def get_copy_and_create_file2__0f2ab16ddcc8f1a930b6b27f18088a10_qw35sft2_1fd84deb(env, config: dict):
    """Check file1 in dir1/dir2/dir3 and new file2 in home dir."""
    result = env.controller.run_bash_script('echo dir1:$(test -f /home/user/dir1/file1 && echo yes || echo no)\necho dir2:$(test -f /home/user/dir2/file1 && echo yes || echo no)\necho dir3:$(test -f /home/user/dir3/file1 && echo yes || echo no)\necho file2:$(test -f /home/user/file2 && echo yes || echo no)', timeout=10)
    if not result:
        return {'error': 'command failed'}
    output = result.get('output', '')
    return {'dir1_has_file1': 'dir1:yes' in output, 'dir2_has_file1': 'dir2:yes' in output, 'dir3_has_file1': 'dir3:yes' in output, 'home_has_file2': 'file2:yes' in output}

def get_file_and_dir_permissions__6736312e2a9f6c78caa7e5e2f35a4133_qw35sft2_e7b0794b(env, config: dict):
    """Get permissions of all regular files and subdirectories under testDir."""
    file_result = env.controller.run_bash_script('find /home/user/testDir -type f -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    dir_result = env.controller.run_bash_script('find /home/user/testDir -mindepth 1 -type d -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    file_perms = [p.strip() for p in file_result.get('stdout', '').strip().split('\n') if p.strip()]
    dir_perms = [p.strip() for p in dir_result.get('stdout', '').strip().split('\n') if p.strip()]
    return {'file_permissions': file_perms, 'dir_permissions': dir_perms}

def get_user_audit_files__c74bb543554ec2beed49099c34c9413f_qw35sft2_d16d2df5(env, config: dict):
    """Read ~/user_audit/whoami.txt and ~/user_audit/users.txt for partial-credit scoring."""
    whoami_bytes = env.controller.get_file('/home/user/user_audit/whoami.txt')
    users_bytes = env.controller.get_file('/home/user/user_audit/users.txt')
    whoami_content = whoami_bytes.decode('utf-8', errors='replace').strip() if whoami_bytes else ''
    users_content = users_bytes.decode('utf-8', errors='replace').strip() if users_bytes else ''
    return {'whoami_content': whoami_content, 'users_content': users_content}

def get_file_org_state__8f46312161f8fe4e75948b31954b5be0_qw35sft2_c656985f(env, config: dict):
    """
    Check:
    1. A .tar.gz archive exists in /tmp/test_files/old_files/ and contains >=2 .txt files
       (proving the 30-day-old files were actually compressed into it).
    2. /tmp/test_files/new_files/ contains at least 2 .txt files
       (recently modified files were moved there by the agent).
    """
    script = 'ARCHIVE=$(find /tmp/test_files/old_files -maxdepth 1 -name "*.tar.gz" | head -1); echo "ARCHIVE:${ARCHIVE}"; if [ -n "${ARCHIVE}" ]; then   tar -tzf "${ARCHIVE}" 2>/dev/null | grep -oE "[^/]+\\.txt$"; fi; echo "---"; find /tmp/test_files/new_files -maxdepth 1 -type f -name "*.txt" | sort'
    result = env.controller.run_bash_script(script, timeout=30)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    archive_path = ''
    archive_txt_files = []
    new_file_names = []
    section = 0
    for line in output.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith('ARCHIVE:'):
            archive_path = line[len('ARCHIVE:'):]
        elif line == '---':
            section = 1
        elif section == 0 and line.endswith('.txt'):
            archive_txt_files.append(line)
        elif section == 1:
            new_file_names.append(line.split('/')[-1])
    return {'archive_found': bool(archive_path), 'archive_path': archive_path, 'archive_txt_files': archive_txt_files, 'archive_has_txt_files': len(archive_txt_files) >= 2, 'new_files': new_file_names, 'new_files_count': len(new_file_names)}

def get_subdir_split_permissions__31b8299b7936d108e594f17a62518506_qw35sft2_a0f0f736(env, config: dict):
    """Get permissions of files in subDir1 and subDir2 separately."""
    subdir1_result = env.controller.run_bash_script('find /home/user/testDir/subDir1 -type f -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    subdir2_result = env.controller.run_bash_script('find /home/user/testDir/subDir2 -type f -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    subdir1_perms = [p.strip() for p in subdir1_result.get('stdout', '').strip().split('\n') if p.strip()]
    subdir2_perms = [p.strip() for p in subdir2_result.get('stdout', '').strip().split('\n') if p.strip()]
    return {'subdir1_permissions': subdir1_perms, 'subdir2_permissions': subdir2_perms}

def get_file_copy_3dirs__d110c9f1795d3fe99db0ae18f13e5b66_qw35sft2_3a1d9288(env, config: dict):
    """Check whether file1 was copied into dir1, dir2, and dir3."""
    result = env.controller.run_bash_script('echo dir1:$(test -f /home/user/dir1/file1 && echo yes || echo no)\necho dir2:$(test -f /home/user/dir2/file1 && echo yes || echo no)\necho dir3:$(test -f /home/user/dir3/file1 && echo yes || echo no)', timeout=10)
    if not result:
        return {'error': 'command failed'}
    output = result.get('output', '')
    return {'dir1_has_file1': 'dir1:yes' in output, 'dir2_has_file1': 'dir2:yes' in output, 'dir3_has_file1': 'dir3:yes' in output}

def get_desktop_files_check__ee992b5c4f2de4a7b0ef0f22ec1a67b1_qw35sft2_e9092f55(env, config: dict):
    """Check Desktop for both old and new filename of the restored poster."""
    vm_ip = env.vm_ip
    port = env.server_port
    try:
        resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'ls /home/user/Desktop/ 2>/dev/null', 'shell': True}, timeout=15)
        if resp.status_code != 200:
            return {'error': f'HTTP {resp.status_code}'}
        output = resp.json().get('output', '') or ''
        files = [f.strip() for f in output.strip().split('\n') if f.strip()]
        return {'has_new_name': 'party_poster.webp' in files, 'has_old_name': 'poster_party_night.webp' in files}
    except Exception as e:
        logger_qw35sft2_a25b79.error('get_desktop_files_check__ee992b5c4f2de4a7b0ef0f22ec1a67b1 error: %s', e)
        return {'error': str(e)}

def get_rename_and_subfolder__cc8a90beff8b6d7aa9d9f79938c158e4_qw35sft2_42301119(env, config: dict):
    """Check if Desktop/todo_list_Jan_2 exists and contains a 'tasks' subfolder."""
    vm_ip = env.vm_ip
    port = env.server_port
    r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_2 && echo yes || echo no', 'shell': True})
    renamed = r1.status_code == 200 and r1.json().get('output', '').strip() == 'yes'
    r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_2/tasks && echo yes || echo no', 'shell': True})
    subfolder = r2.status_code == 200 and r2.json().get('output', '').strip() == 'yes'
    return {'renamed': renamed, 'subfolder': subfolder}

def get_new_files_perms__2e6f4a565bdf8ce8ae9bb4e56e1c79f7_qw35sft2_218fed0a(env, config: dict):
    """Check archive exists in old_files, recent files moved to new_files, and permissions of new_files directory."""
    result = env.controller.run_bash_script('ARCH=$(find /tmp/test_files/old_files -maxdepth 1 -name "*.tar.gz" | head -1); PERM=$(stat -c "%a" /tmp/test_files/new_files 2>/dev/null); COUNT=$(find /tmp/test_files/new_files -maxdepth 1 -type f 2>/dev/null | wc -l); echo "ARCHIVE=${ARCH}"; echo "PERM=${PERM}"; echo "COUNT=${COUNT}"', timeout=15)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    archive_path = ''
    perm = ''
    count_str = '0'
    for line in output.strip().splitlines():
        if line.startswith('ARCHIVE='):
            archive_path = line[len('ARCHIVE='):].strip()
        elif line.startswith('PERM='):
            perm = line[len('PERM='):].strip()
        elif line.startswith('COUNT='):
            count_str = line[len('COUNT='):].strip()
    try:
        file_count = int(count_str)
    except ValueError:
        file_count = 0
    return {'archive_exists': bool(archive_path), 'new_files_populated': file_count > 0, 'new_files_perm': perm}

def get_home_users_file__baeeb37009528a1d856ed9e685fd5d89_qw35sft2_de61165a(env, config: dict):
    """Read ~/home_users.txt containing /etc/passwd entries for users with /home/ dirs."""
    path = '/home/user/home_users.txt'
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    return {'content': content}

def get_two_file_state__bb3776631f95d3c4e2bcd96582b451ff_qw35sft2_a5bedd98(env, config: dict):
    """Get content of output.txt and count.txt from the VM."""
    vm_ip = env.vm_ip
    port = env.server_port
    output_content = ''
    resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['cat', '/home/user/output.txt'], 'shell': False})
    if resp.status_code == 200:
        output_content = resp.json().get('output', '')
    count_content = ''
    resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['cat', '/home/user/count.txt'], 'shell': False})
    if resp.status_code == 200:
        count_content = resp.json().get('output', '')
    return {'output_txt': output_content, 'count_txt': count_content}

def get_rename_and_file__caa3a14501f9d1aa7da658fbd50adb72_qw35sft2_ca1c5ca0(env, config: dict):
    """Check if Desktop/todo_list_Jan_2 exists and contains notes.txt."""
    vm_ip = env.vm_ip
    port = env.server_port
    r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_2 && echo yes || echo no', 'shell': True})
    renamed = r1.status_code == 200 and r1.json().get('output', '').strip() == 'yes'
    r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -f ~/Desktop/todo_list_Jan_2/notes.txt && echo yes || echo no', 'shell': True})
    file_exists = r2.status_code == 200 and r2.json().get('output', '').strip() == 'yes'
    return {'renamed': renamed, 'file_exists': file_exists}

def get_file_perms_and_archive__00b9eddcbcb085dc89e5e17c3b839024_qw35sft2_9bf6c2cc(env, config: dict):
    """Get permissions of all regular files and check for testDir.tar archive."""
    perm_result = env.controller.run_bash_script('find /home/user/testDir -type f -exec stat -c "%a" {} \\; 2>/dev/null', timeout=15)
    archive_result = env.controller.run_bash_script('test -f /home/user/testDir.tar && echo "exists" || echo "missing"', timeout=10)
    file_perms = [p.strip() for p in perm_result.get('stdout', '').strip().split('\n') if p.strip()]
    archive_exists = archive_result.get('stdout', '').strip() == 'exists'
    return {'file_permissions': file_perms, 'archive_exists': archive_exists}

def get_archive_content__7b7da38700f679a895053728c5a7d35b_qw35sft2_1ee11bd6(env, config: dict):
    """Find the tar.gz archive in old_files and list its contents."""
    result = env.controller.run_bash_script('ARCH=$(find /tmp/test_files/old_files -maxdepth 1 -name "*.tar.gz" | head -1); if [ -n "$ARCH" ]; then echo "ARCHIVE_EXISTS"; tar -tzf "$ARCH" 2>/dev/null; else echo "NO_ARCHIVE"; fi', timeout=15)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    lines = [l.strip() for l in output.strip().splitlines() if l.strip()]
    archive_exists = lines[0] == 'ARCHIVE_EXISTS' if lines else False
    contents = []
    if archive_exists and len(lines) > 1:
        contents = [l.split('/')[-1] for l in lines[1:] if l]
    return {'archive_exists': archive_exists, 'contents': contents}

def get_two_dir_notebook_split__e71563318b46b13cadfcb76e6b9ce246_qw35sft2_71e98bf9(env, config: dict):
    """Get files in ./fails (failed notebooks) and ./passing (non-failed notebooks)."""
    vm_ip = env.vm_ip
    port = env.server_port
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'find /home/user/test_environment/fails -type f 2>/dev/null | sort', 'shell': True})
    fails_files = resp1.json().get('output', '').strip() if resp1.status_code == 200 else ''
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'find /home/user/test_environment/passing -type f 2>/dev/null | sort', 'shell': True})
    passing_files = resp2.json().get('output', '').strip() if resp2.status_code == 200 else ''
    return {'fails_files': fails_files, 'passing_files': passing_files}

def get_bills_and_local_folders__0b7b1b3c91e02a8222608f02952fc214_qw35sft2_8ba964c8(env, config: dict):
    """Read Bills mbox starred status AND check for Receipts folder existence."""
    mbox_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills'
    file_bytes = env.controller.get_file(mbox_path)
    all_starred = False
    message_count = 0
    if file_bytes:
        content = file_bytes.decode('utf-8', errors='replace')
        messages = []
        parts = re.split('^From - ', content, flags=re.MULTILINE)
        for part in parts:
            if not part.strip():
                continue
            status_match = re.search('^X-Mozilla-Status: ([0-9A-Fa-f]+)', part, re.MULTILINE)
            if not status_match:
                continue
            status_int = int(status_match.group(1), 16)
            if status_int & 8:
                continue
            messages.append(bool(status_int & 4))
        message_count = len(messages)
        all_starred = bool(messages) and all(messages)
    receipts_msf = env.controller.get_file('/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Receipts.msf')
    receipts_exists = receipts_msf is not None
    return {'message_count': message_count, 'all_starred': all_starred, 'receipts_folder_exists': receipts_exists}

def get_file_exists__df188ad33aeda3e315e62cc2c6afb173_qw35sft2_f28ae247(env, config: dict):
    """Check whether a file exists at the specified VM path."""
    path = config.get('path', '')
    if not path:
        return {'exists': False, 'error': 'no path provided'}
    result = env.controller.run_bash_script(f'test -f "{path}" && echo "exists" || echo "not_found"', timeout=10)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    return {'exists': 'exists' in output, 'path': path}

def get_main_py_content__2c85ba49c59e818794c6aa64361a93f2_qw35sft2_47dcf891(env, config: dict):
    """Read /home/user/main.py from the VM and return its text content."""
    file_bytes = env.controller.get_file('/home/user/main.py')
    if not file_bytes:
        return {'error': 'File not found'}
    return {'content': file_bytes.decode('utf-8', errors='replace')}

def get_main_py_content__4c84b586b63f7c1cf9b6df39aedcd5ed_qw35sft2_db7ad539(env, config: dict):
    """Read /home/user/main.py from the VM and return its text content."""
    file_bytes = env.controller.get_file('/home/user/main.py')
    if not file_bytes:
        return {'error': 'File not found'}
    return {'content': file_bytes.decode('utf-8', errors='replace')}
