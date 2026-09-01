"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_html_file_content__43b4c11e2d677db2c4e2a819d325a693', 'get_bookmark_folder_with_urls__1eba9f123a42adb66832822e9cf027cd', 'get_docx_table_text__9467a3bac5aae80d6f5894cff6144c37', 'get_html_file_content__0312c898008bbe83c35bf8d3f0838a0a', 'get_docx_table_info__126b20975e4e6ed136cbb81488fd0810', 'get_active_url_from_accessTree', 'get_docx_table_text__75b704f601003cb29b05bac44a154402', 'get_browser_active_url__a5eff305d3fedc2c22bac974fb8698d7', 'get_bookmark_and_active_tab__a0bd39d93c2a735cb19f1b8174d16ac1', 'get_docx_table_text__257a7a045a00e35dcfe1e7d02a439bcc', 'get_chrome_active_url__9f2f3fb3a73693af5f879fd7c403da86', 'get_chrome_startup_setting__8434e9ea78e4611ce770a316aad99fa9', 'get_docx_table_bold__b70d187abe0811a915df187857b97f14', 'get_chrome_color_scheme__c109c66a806aa3929078e7fafeda76b2_qw35sft2_d1928ab1', 'get_chrome_startup_and_font__aa9af3d49c1447b85e91f3d729dc4ea1_qw35sft2_c2f0067b', 'get_chrome_delete_and_safebrowsing__be8c0e2643b53348f7e3a53d09b2fd8d_qw35sft2_684d1d9d', 'get_chrome_dnt_and_lang__6c4389c0ed5d5685554268779bf7b6e2_qw35sft2_a0e5c9e2', 'get_chrome_profile_bookmark__f302c5cc6d17b170f6df5e52c54310d0_qw35sft2_96d8156c', 'get_chrome_extensions_dev_mode__b759dda996122d64d46867fb7e10f84a_qw35sft2_93149557', 'get_third_party_cookie_mode__2e738e6837968a7fb46f2d87cbf75561_qw35sft2_a477ed9e', 'get_chrome_search_bookmarks__7963cff65d97676543d55b7ec9b158c6_qw35sft2_5e2dc99e', 'get_chrome_sb_and_dnt__497c76e8a817211059a0ca1d7eafcbe5_qw35sft2_845d10b2', 'get_chrome_dnt_and_font__f0db2c12ec96d3956f7c7b7b8c915fc1_qw35sft2_c8515ee5', 'get_chrome_profile_search__bae1be09edcf3b3e43698d058f3c0784_qw35sft2_24cb55ab', 'get_chrome_startup_urls__db8335a9df035ba863791bb757e40ea6_qw35sft2_14cce345', 'get_chrome_delete_and_startup__f2a35c2c45e7134d422e18016f6e0ff7_qw35sft2_8ef80023', 'get_chrome_unpacked_extensions__83ea9c43b50b1e067c01ffbea0552893_qw35sft2_8f406767', 'get_google_search_toggle__d5be297d19eda82733e6c2f634aba9e2_qw35sft2_4a54b6d4', 'get_bookmark_folder_and_dnt__33cb989983633f67656c84eb11865091_qw35sft2_9957cbaa', 'get_chrome_appearance_and_dnt__197c5f21c248c0028a67e57d9193addd_qw35sft2_0f0a9cca', 'get_chrome_search_dnt_state__67085c82c40b9ce759bdc0050139fed9_qw35sft2_10b9c364', 'get_chrome_lang_and_dnt__45fdf66a8d514412f8287b18b125d4f9_qw35sft2_14044afb', 'get_bookmarks_bar_urls__34222c255c19ca63c5f6efa3ee3ba731_qw35sft2_d29a6ea1', 'get_chrome_delete_and_dnt__147c27e083ffc1ed15c90d12c99f1fef_qw35sft2_e3c218a8', 'get_darktable_installed__bbd0b6cc65b6ed22ef194ee8993ef563_qw35sft2_022411c2', 'get_docx_table_first_row__a4c465e58afc03255b4d0ddc46d5d525_qw35sft2_0ce40d24', 'get_doc_tabstops__5d3adaf927fd3613be8530bc92e14de1_qw35sft2_e618f7eb', 'get_docx_table_italic_col__838d41f5c432af14f1eb478fa956e99e_qw35sft2_46821bef', 'get_planet_table_rows__5135726bf1592f829eb43617d4e867fa_qw35sft2_f66e74b7', 'get_doc_tabstops__14c984754ad8a9510e10b62f00a4c316_qw35sft2_e918c2d3', 'get_docx_table_structure__638d8a63e2ebfdf518128094d93e23f6_qw35sft2_5c88a236', 'get_doc_tabstops__319c26794056ac5ca6a78fe65ba6a022_qw35sft2_e7ab72be', 'get_docx_table_header_align__83874f16b468982748e51b1480830a3a_qw35sft2_37bedf11', 'get_docx_table_dims__94b525d2310c47dd52401ee93f279d0a_qw35sft2_2214bd9c', 'get_doc_tabstops__cf559754b29c29c8081fb54264f1ac10_qw35sft2_23c444b3', 'get_docx_table_bold_header__d0985f126b123d58872633f0350725ab_qw35sft2_3a1d6b53', 'get_docx_table_dims__f2fed6259df9a3cde848f19734b96eb2_qw35sft2_ac504bf7', 'get_docx_table_last_row__e0b81f05c0001516085ceb6f32f7b48c_qw35sft2_b0c3fdf7', 'get_timetable_remove_tutorial6__c0f66209d7bea00756990ac3a15deab7_qw35sft2_f9a85fc1', 'get_chrome_tabs_and_bookmarks__ba148d88a9035c4354cf728e6417abab_qw35sft2_bb9b51d7', 'get_timetable_lec_color__462e6cf4d0037b83d0687915bbffba64_qw35sft2_157dbc7f', 'get_extension_project_state__ce07ffbd7d3847464df6caec281dec8d_qw35sft2_770a4b6d']

def get_html_file_content__43b4c11e2d677db2c4e2a819d325a693(env, config: dict):
    """Get HTML file content from VM and check existence."""
    import re
    path = config.get('path', '')
    if not path:
        return {'error': 'No path specified'}
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'exists': False, 'content': '', 'title': ''}
    except Exception:
        return {'exists': False, 'content': '', 'title': ''}
    content = file_bytes.decode('utf-8', errors='replace')
    title_match = re.search('<title>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)
    title = title_match.group(1).strip() if title_match else ''
    return {'exists': True, 'content': content, 'title': title}

def get_bookmark_folder_with_urls__1eba9f123a42adb66832822e9cf027cd(env, config: dict):
    """Get bookmark folder info and active tab URL for partial credit evaluation."""
    result = {}
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Bookmarks'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Bookmarks'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Bookmarks'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Bookmarks'))")['output'].strip()
    else:
        return {'error': 'Unsupported OS'}
    try:
        bookmarks_content = env.controller.get_file(preference_file_path)
        if bookmarks_content:
            bookmarks = json.loads(bookmarks_content)
            bookmark_bar = bookmarks.get('roots', {}).get('bookmark_bar', {})
            children = bookmark_bar.get('children', [])
            target_folder = None
            for child in children:
                if child.get('type') == 'folder' and child.get('name') == config.get('folder_name', 'Papers'):
                    target_folder = child
                    break
            if target_folder:
                result['folder_exists'] = True
                folder_urls = [c.get('url', '') for c in target_folder.get('children', []) if c.get('type') == 'url']
                result['folder_urls'] = folder_urls
            else:
                result['folder_exists'] = False
                result['folder_urls'] = []
        else:
            result['folder_exists'] = False
            result['folder_urls'] = []
    except Exception as e:
        logger.error(f'Error reading bookmarks: {e}')
        result['folder_exists'] = False
        result['folder_urls'] = []
    return result

def get_docx_table_text__9467a3bac5aae80d6f5894cff6144c37(env, config: dict):
    """Get all table data from docx to check for GPT-4 specific scores."""
    from docx import Document
    file_path = config.get('path', '/home/user/Documents/awesome-desktop/awe_desk_env.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'tables': [], 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        for (i, table) in enumerate(doc.tables):
            rows_data = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_text)
            tables_data.append(rows_data)
        return {'tables': tables_data, 'table_count': len(tables_data)}
    finally:
        os.unlink(tmp_path)

def get_html_file_content__0312c898008bbe83c35bf8d3f0838a0a(env, config: dict):
    """Get HTML file content from VM and return existence, size, and content."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    content = file_bytes.decode('utf-8', errors='ignore')
    return {'exists': True, 'size': len(file_bytes), 'content_lower': content.lower()}

def get_docx_table_info__126b20975e4e6ed136cbb81488fd0810(env, config: dict):
    """Get table count and dimensions of the last table in a docx file."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_table_rows': 0, 'last_table_cols': 0}
        last_table = tables[-1]
        rows = len(last_table.rows)
        cols = len(last_table.columns)
        return {'table_count': table_count, 'last_table_rows': rows, 'last_table_cols': cols}
    finally:
        os.unlink(tmp_path)

def get_active_url_from_accessTree(env, config: Dict[str, str]):
    """Get active URL from accessibility tree"""
    import lxml.etree
    _accessibility_ns_map = {'st': 'uri:deskat:state.at-spi.gnome.org', 'attr': 'uri:deskat:attributes.at-spi.gnome.org', 'cp': 'uri:deskat:component.at-spi.gnome.org', 'doc': 'uri:deskat:document.at-spi.gnome.org', 'docattr': 'uri:deskat:attributes.document.at-spi.gnome.org', 'txt': 'uri:deskat:text.at-spi.gnome.org', 'val': 'uri:deskat:value.at-spi.gnome.org', 'act': 'uri:deskat:action.at-spi.gnome.org'}
    xml_str = env.controller.get_accessibility_tree()
    if xml_str is None:
        return None
    root = lxml.etree.fromstring(xml_str.encode('utf-8'))
    prefix = config.get('goto_prefix', 'https://www.')
    for frame in root.xpath("//frame[@st:focused='true']", namespaces=_accessibility_ns_map):
        name = frame.get('name')
        if name and name.startswith(prefix):
            return name
    return None

def get_docx_table_text__75b704f601003cb29b05bac44a154402(env, config: dict):
    """Get all table text from docx file to check if a table was inserted with GPT-4 name and avg score."""
    from docx import Document
    file_path = config.get('path', '/home/user/Documents/awesome-desktop/awe_desk_env.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'tables': [], 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        for (i, table) in enumerate(doc.tables):
            rows_data = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_text)
            tables_data.append(rows_data)
        return {'tables': tables_data, 'table_count': len(tables_data)}
    finally:
        os.unlink(tmp_path)

def get_browser_active_url__a5eff305d3fedc2c22bac974fb8698d7(env, config: dict):
    """Get the current active URL from the browser using Chrome DevTools Protocol."""
    try:
        result = env.controller.run_bash_script('curl -s http://localhost:9222/json', timeout=10)
        if isinstance(result, dict):
            output = result.get('output', '')
        else:
            output = str(result)
        if not output.strip():
            return {'error': 'No response from CDP endpoint'}
        tabs = json.loads(output)
        if not tabs:
            return {'error': 'No tabs found'}
        active_url = tabs[0].get('url', '')
        return {'url': active_url}
    except json.JSONDecodeError as e:
        return {'error': f'Failed to parse CDP response: {str(e)}'}
    except Exception as e:
        return {'error': str(e)}

def get_bookmark_and_active_tab__a0bd39d93c2a735cb19f1b8174d16ac1(env, config: dict):
    """Get bookmarks bar URLs and the active tab URL."""
    import json
    bookmarks_raw = env.controller.get_bookmarks()
    bookmark_urls = []
    if isinstance(bookmarks_raw, dict):
        bar = bookmarks_raw.get('bookmark_bar', {})
        children = bar.get('children', [])
        for child in children:
            if 'url' in child:
                bookmark_urls.append(child['url'])
    tree = env.controller.get_accessibility_tree()
    active_url = ''
    if isinstance(tree, str):
        for line in tree.split('\n'):
            if 'document' in line.lower() and 'url=' in line.lower():
                idx = line.lower().find('url=')
                if idx >= 0:
                    active_url = line[idx + 4:].strip().strip('\'"')
                    break
    if not active_url:
        try:
            tabs = env.controller.get_open_tabs_info()
            if tabs and len(tabs) > 0:
                for tab in tabs:
                    if isinstance(tab, dict) and tab.get('active', False):
                        active_url = tab.get('url', '')
                        break
                if not active_url and isinstance(tabs[-1], dict):
                    active_url = tabs[-1].get('url', '')
        except Exception:
            pass
    return {'bookmark_urls': bookmark_urls, 'active_url': active_url}

def get_docx_table_text__257a7a045a00e35dcfe1e7d02a439bcc(env, config: dict):
    """Get all table text from docx file to check if a table was inserted with model names."""
    from docx import Document
    file_path = config.get('path', '/home/user/Documents/awesome-desktop/awe_desk_env.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'tables': [], 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        for (i, table) in enumerate(doc.tables):
            rows_data = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                rows_data.append(row_text)
            tables_data.append(rows_data)
        section_text = ''
        in_section = False
        for para in doc.paragraphs:
            text = para.text.strip()
            if 'Main Results' in text:
                in_section = True
                continue
            if in_section and ('Qualitative Analysis' in text or 'Conclusion' in text):
                in_section = False
                break
            if in_section:
                section_text += text + ' '
        return {'tables': tables_data, 'table_count': len(tables_data), 'main_results_text': section_text.strip()}
    finally:
        os.unlink(tmp_path)

def get_chrome_active_url__9f2f3fb3a73693af5f879fd7c403da86(env, config: dict):
    """Get the active tab URL from Chrome via accessibility tree."""
    try:
        tree = env.controller.get_accessibility_tree()
        if not tree:
            return {'error': 'Could not get accessibility tree'}
        import re
        for line in tree.split('\n'):
            if 'address' in line.lower() and 'bar' in line.lower():
                url_match = re.search('(https?://[^\\s"\\\']+)', line)
                if url_match:
                    return {'url': url_match.group(1)}
            if 'taoyds.github.io' in line or 'leuchine.github.io' in line or 'ikekonglp.github.io' in line:
                url_match = re.search('(https?://[^\\s"\\\']+)', line)
                if url_match:
                    return {'url': url_match.group(1)}
        result = env.controller.run_bash_script('python3 -c "import subprocess, json;import urllib.request;req = urllib.request.Request(\'http://localhost:9222/json\');resp = urllib.request.urlopen(req);tabs = json.loads(resp.read());print(tabs[0][\'url\'] if tabs else \'\')"', timeout=10)
        if result and isinstance(result, dict):
            output = result.get('output', '').strip()
            if output:
                return {'url': output}
        return {'error': 'Could not extract URL'}
    except Exception as e:
        logger.error(f'Error getting Chrome URL: {e}')
        return {'error': str(e)}

def get_chrome_startup_setting__8434e9ea78e4611ce770a316aad99fa9(env, config: dict):
    """Get Chrome's startup setting (restore_on_startup) from Preferences file."""
    import json
    try:
        file_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
        if not file_bytes:
            return {'error': 'Preferences file not found'}
        prefs = json.loads(file_bytes.decode('utf-8'))
        session = prefs.get('session', {})
        return {'restore_on_startup': session.get('restore_on_startup', -1), 'startup_urls': session.get('startup_urls', [])}
    except Exception as e:
        return {'error': str(e)}

def get_docx_table_bold__b70d187abe0811a915df187857b97f14(env, config: dict):
    """Get bold status of all words in the document's table."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        total_words = 0
        bold_words = 0
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                total_words += 1
                                if run.font.bold:
                                    bold_words += 1
        return {'total_words': total_words, 'bold_words': bold_words, 'all_bold': total_words > 0 and bold_words == total_words}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_chrome_color_scheme__c109c66a806aa3929078e7fafeda76b2_qw35sft2_d1928ab1(env, config: dict):
    """Get Chrome color scheme mode as a string: 'system', 'light', or 'dark'."""
    prefs_path = '/home/user/.config/google-chrome/Default/Preferences'
    SCHEME_MAP = {0: 'system', 1: 'light', 2: 'dark'}
    try:
        content = env.controller.get_file(prefs_path)
        if not content:
            return 'unknown'
        data = json.loads(content)
        color_scheme = data.get('browser', {}).get('theme', {}).get('color_scheme', -1)
        return SCHEME_MAP.get(color_scheme, 'unknown')
    except Exception:
        return 'unknown'

def get_chrome_startup_and_font__aa9af3d49c1447b85e91f3d729dc4ea1_qw35sft2_c2f0067b(env, config: dict):
    """
    Get Chrome startup config and font size settings from Preferences file.
    Returns dict with: restore_on_startup (int), startup_urls (list), default_font_size (int).
    """
    preference_file_path = '/home/user/.config/google-chrome/Default/Preferences'
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        session = data.get('session', {})
        restore_on_startup = session.get('restore_on_startup', -1)
        startup_urls = session.get('startup_urls', [])
        font_size = data.get('webkit', {}).get('webprefs', {}).get('default_font_size', 16)
        return {'restore_on_startup': restore_on_startup, 'startup_urls': startup_urls if isinstance(startup_urls, list) else [], 'default_font_size': font_size}
    except Exception:
        return {'restore_on_startup': -1, 'startup_urls': [], 'default_font_size': 16}

def get_chrome_delete_and_safebrowsing__be8c0e2643b53348f7e3a53d09b2fd8d_qw35sft2_684d1d9d(env, config: dict):
    """Get combined Chrome state: auto-delete site data on close + Enhanced Safe Browsing."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    else:
        return {'error': 'Unsupported OS'}
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        cookies_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('cookies')
        enhanced_safe_browsing = data.get('safebrowsing', {}).get('enhanced', False)
        return {'data_delete': 'true' if cookies_setting == 4 else 'false', 'enhanced_safe_browsing': 'true' if enhanced_safe_browsing else 'false'}
    except Exception as e:
        return {'error': str(e)}

def get_chrome_dnt_and_lang__6c4389c0ed5d5685554268779bf7b6e2_qw35sft2_a0e5c9e2(env, config: dict):
    """Get Chrome DNT status and interface language from preferences file."""
    import json
    import tempfile
    import os
    file_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
    if not file_bytes:
        return {'error': 'Preferences file not found', 'dnt': False, 'language': None}
    with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            prefs = json.load(f)
        dnt = bool(prefs.get('enable_do_not_track', False))
        language = prefs.get('intl', {}).get('app_locale', '')
        return {'dnt': dnt, 'language': language}
    except Exception as e:
        return {'error': str(e), 'dnt': False, 'language': None}
    finally:
        os.unlink(tmp_path)

def get_chrome_profile_bookmark__f302c5cc6d17b170f6df5e52c54310d0_qw35sft2_96d8156c(env, config: dict):
    """Get Chrome profile name and extract all bookmark URLs from the bookmarks bar."""
    import json
    import tempfile
    import os
    pref_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
    bkm_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Bookmarks')
    profile_name = None
    if pref_bytes:
        try:
            with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
                tmp.write(pref_bytes)
                tmp_path = tmp.name
            try:
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    prefs = json.load(f)
                profile_name = prefs.get('profile', {}).get('name', None)
            finally:
                os.unlink(tmp_path)
        except Exception:
            pass
    bookmark_bar_urls = []
    if bkm_bytes:
        try:
            with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
                tmp.write(bkm_bytes)
                tmp_path = tmp.name
            try:
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    bkm_data = json.load(f)
                roots = bkm_data.get('roots', {})
                stack = [roots.get('bookmark_bar', {})]
                while stack:
                    node = stack.pop()
                    if isinstance(node, dict):
                        if node.get('type') == 'url':
                            url = node.get('url', '')
                            if url:
                                bookmark_bar_urls.append(url)
                        stack.extend(node.get('children', []))
                    elif isinstance(node, list):
                        stack.extend(node)
            finally:
                os.unlink(tmp_path)
        except Exception:
            pass
    return {'profile_name': profile_name, 'bookmark_bar_urls': bookmark_bar_urls}

def get_chrome_extensions_dev_mode__b759dda996122d64d46867fb7e10f84a_qw35sft2_93149557(env, config: dict):
    """Check if Developer mode is enabled in Chrome extensions via preferences file."""
    result = env.controller.run_bash_script('python3 -c "import json; prefs = json.load(open(\'/home/user/.config/google-chrome/Default/Preferences\')); val = prefs.get(\'extensions\', {}).get(\'ui\', {}).get(\'developer_mode\', False); print(bool(val))"', timeout=15)
    if isinstance(result, dict) and result.get('output'):
        out = result['output'].strip()
        return {'developer_mode': out == 'True'}
    return {'developer_mode': False}

def get_third_party_cookie_mode__2e738e6837968a7fb46f2d87cbf75561_qw35sft2_a477ed9e(env, config: dict):
    """Read Chrome's third-party cookie blocking mode from Preferences.

    Returns:
        dict with 'cookie_controls_mode':
            0 = Allow all third-party cookies
            1 = Block third-party cookies in Incognito only
            2 = Block all third-party cookies
    """
    import json
    result = env.controller.run_bash_script('cat /home/user/.config/google-chrome/Default/Preferences 2>/dev/null', timeout=15)
    try:
        raw = result.get('output', '') if isinstance(result, dict) else str(result)
        prefs = json.loads(raw)
        mode = prefs.get('profile', {}).get('cookie_controls_mode', 0)
        return {'cookie_controls_mode': int(mode)}
    except Exception as exc:
        return {'error': str(exc), 'cookie_controls_mode': -1}

def get_chrome_search_bookmarks__7963cff65d97676543d55b7ec9b158c6_qw35sft2_5e2dc99e(env, config: dict):
    """Get Chrome default search engine name and whether bing.com is in the bookmarks bar."""
    import json
    import platform
    os_type = env.vm_platform
    if os_type == 'Linux':
        if 'arm' in platform.machine():
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
            bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Bookmarks'))")['output'].strip()
        else:
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
            bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Bookmarks'))")['output'].strip()
    else:
        prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
        bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Bookmarks'))")['output'].strip()
    try:
        content = env.controller.get_file(prefs_path)
        data = json.loads(content)
        search_engine = data.get('default_search_provider_data', {}).get('template_url_data', {}).get('short_name', 'Google')
    except Exception:
        search_engine = 'Google'
    try:
        bm_content = env.controller.get_file(bm_path)
        bm_data = json.loads(bm_content)
        bookmark_bar = bm_data.get('roots', {}).get('bookmark_bar', {})
        bing_bookmarked = False
        stack = [bookmark_bar]
        while stack:
            node = stack.pop()
            if node.get('type') == 'url' and 'bing.com' in node.get('url', ''):
                bing_bookmarked = True
                break
            stack.extend(node.get('children', []))
        return {'search_engine': search_engine, 'bing_bookmarked': bing_bookmarked}
    except Exception as e:
        return {'search_engine': search_engine, 'bing_bookmarked': False, 'error': str(e)}

def get_chrome_sb_and_dnt__497c76e8a817211059a0ca1d7eafcbe5_qw35sft2_845d10b2(env, config: dict):
    """Get Chrome Safe Browsing and Do Not Track settings from Preferences.

    Returns dict:
      safe_browsing: "true" if standard or enhanced protection is on, "false" otherwise
      do_not_track: "true" if Do Not Track is enabled, "false" otherwise
    """
    pref_path = '/home/user/.config/google-chrome/Default/Preferences'
    try:
        content = env.controller.get_file(pref_path)
        data = json.loads(content)
        sb = data.get('safebrowsing', {})
        safe_browsing = bool(sb.get('enhanced', False) or sb.get('enabled', False))
        do_not_track = bool(data.get('enable_do_not_track', False))
        return {'safe_browsing': 'true' if safe_browsing else 'false', 'do_not_track': 'true' if do_not_track else 'false'}
    except Exception as e:
        logger_qw35sft2_a8fd39.error(f'Error reading Chrome preferences: {e}')
        return {'safe_browsing': 'false', 'do_not_track': 'false'}

def get_chrome_dnt_and_font__f0db2c12ec96d3956f7c7b7b8c915fc1_qw35sft2_c8515ee5(env, config: dict):
    """Get Chrome DNT status and default font size from preferences file."""
    import json
    import tempfile
    import os
    file_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
    if not file_bytes:
        return {'error': 'Preferences file not found', 'dnt': False, 'font_size': None}
    with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            prefs = json.load(f)
        dnt = bool(prefs.get('enable_do_not_track', False))
        font_size = prefs.get('webkit', {}).get('webprefs', {}).get('default_font_size', None)
        return {'dnt': dnt, 'font_size': font_size}
    except Exception as e:
        return {'error': str(e), 'dnt': False, 'font_size': None}
    finally:
        os.unlink(tmp_path)

def get_chrome_profile_search__bae1be09edcf3b3e43698d058f3c0784_qw35sft2_24cb55ab(env, config: dict):
    """Get Chrome profile name and default search engine from preferences file."""
    import json
    import tempfile
    import os
    file_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
    if not file_bytes:
        return {'error': 'Preferences file not found', 'profile_name': None, 'search_engine': None}
    with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            prefs = json.load(f)
        profile_name = prefs.get('profile', {}).get('name', None)
        search_engine = prefs.get('default_search_provider_data', {}).get('template_url_data', {}).get('short_name', 'Google')
        return {'profile_name': profile_name, 'search_engine': search_engine}
    except Exception as e:
        return {'error': str(e), 'profile_name': None, 'search_engine': None}
    finally:
        os.unlink(tmp_path)

def get_chrome_startup_urls__db8335a9df035ba863791bb757e40ea6_qw35sft2_14cce345(env, config: dict):
    """
    Get Chrome startup URLs and mode from Preferences file.
    Returns dict with: restore_on_startup (int), startup_urls (list).
    restore_on_startup values: 1=last session, 4=specific pages, 5=new tab page
    """
    preference_file_path = '/home/user/.config/google-chrome/Default/Preferences'
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        session = data.get('session', {})
        restore_on_startup = session.get('restore_on_startup', -1)
        startup_urls = session.get('startup_urls', [])
        return {'restore_on_startup': restore_on_startup, 'startup_urls': startup_urls if isinstance(startup_urls, list) else []}
    except Exception:
        return {'restore_on_startup': -1, 'startup_urls': []}

def get_chrome_delete_and_startup__f2a35c2c45e7134d422e18016f6e0ff7_qw35sft2_8ef80023(env, config: dict):
    """Get combined Chrome state: auto-delete site data on close + continue where left off on startup."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    else:
        return {'error': 'Unsupported OS'}
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        data_delete_enabled = data.get('profile', {}).get('default_content_setting_values', {}).get('cookies') == 4
        restore_on_startup_ok = data.get('session', {}).get('restore_on_startup', 0) == 5
        return {'data_delete': 'true' if data_delete_enabled else 'false', 'restore_on_startup': 'true' if restore_on_startup_ok else 'false'}
    except Exception as e:
        return {'error': str(e)}

def get_chrome_unpacked_extensions__83ea9c43b50b1e067c01ffbea0552893_qw35sft2_8f406767(env, config: dict):
    """Get names of developer-loaded unpacked extensions from Chrome preferences."""
    script = 'python3 -c "import json; prefs = json.load(open(\'/home/user/.config/google-chrome/Default/Preferences\')); settings = prefs.get(\'extensions\', {}).get(\'settings\', {}); names = [v.get(\'manifest\', {}).get(\'name\', \'\') for v in settings.values() if v.get(\'location\', 0) == 4 and v.get(\'manifest\')]; print(names)"'
    result = env.controller.run_bash_script(script, timeout=15)
    if isinstance(result, dict) and result.get('output'):
        output = result['output'].strip()
        try:
            import ast
            names = ast.literal_eval(output)
            return {'unpacked_extensions': names}
        except Exception:
            return {'unpacked_extensions': []}
    return {'unpacked_extensions': []}

def get_google_search_toggle__d5be297d19eda82733e6c2f634aba9e2_qw35sft2_4a54b6d4(env, config: dict):
    """Get the aria-checked state of a toggle on Google Search Settings via Chrome CDP.

    config keys:
        toggle_label (str): visible text label of the toggle

    Returns dict with 'checked' (True/False/None) and optional 'error'.
    """
    import json
    import base64
    import re
    toggle_label = config.get('toggle_label', '')
    py_script = 'import json, urllib.request, sys\ntoggle_label = ' + json.dumps(toggle_label) + '\ntry:\n    raw = urllib.request.urlopen(\'http://localhost:9222/json\', timeout=5).read()\n    tabs = json.loads(raw)\nexcept Exception as e:\n    print(json.dumps({\'state\': \'cdp_error\', \'error\': str(e)}))\n    sys.exit(0)\nprefs_tab = None\nfor t in tabs:\n    if \'google.com/preferences\' in t.get(\'url\', \'\') and t.get(\'type\') == \'page\':\n        prefs_tab = t\n        break\nif not prefs_tab:\n    for t in reversed(tabs):\n        if t.get(\'type\') == \'page\':\n            prefs_tab = t\n            break\nif not prefs_tab:\n    print(json.dumps({\'state\': \'no_tab\'}))\n    sys.exit(0)\nws_url = prefs_tab.get(\'webSocketDebuggerUrl\', \'\')\ntry:\n    import websocket\n    ws = websocket.create_connection(ws_url, timeout=10)\n    label = toggle_label.replace(\'\\"\', \'\\\\"\').replace("\'", "\\\\\'")\n    js = (\n        \'(function(){\'\n        \'var all=document.querySelectorAll("span,div,label");\'\n        \'for(var i=0;i<all.length;i++){\'\n        \'var el=all[i];\'\n        \'if(el.childElementCount===0&&el.textContent.trim()==="\' + label + \'"){\'\n        \'var p=el.parentElement;\'\n        \'for(var d=0;d<10;d++){\'\n        \'if(!p)break;\'\n        \'var sw=p.querySelector("[role=\\\\"switch\\\\"]");\'\n        \'if(sw)return sw.getAttribute("aria-checked");\'\n        \'p=p.parentElement;\'\n        \'}\'\n        \'}\'\n        \'}\'\n        \'return "not_found";\'\n        \'})()\')\n    ws.send(json.dumps({\'id\': 1, \'method\': \'Runtime.evaluate\',\n                        \'params\': {\'expression\': js, \'returnByValue\': True}}))\n    resp = json.loads(ws.recv())\n    ws.close()\n    val = resp.get(\'result\', {}).get(\'result\', {}).get(\'value\', \'not_found\')\n    print(json.dumps({\'state\': val}))\nexcept ImportError:\n    print(json.dumps({\'state\': \'no_websocket\'}))\nexcept Exception as e:\n    print(json.dumps({\'state\': \'ws_error\', \'error\': str(e)}))\n'
    encoded = base64.b64encode(py_script.encode('utf-8')).decode('ascii')
    write_cmd = 'python3 -c "import base64; open(\'/tmp/gst_v4.py\',\'w\').write(base64.b64decode(\'' + encoded + '\').decode(\'utf-8\'))"'
    env.controller.run_bash_script(write_cmd, timeout=10)
    result = env.controller.run_bash_script('python3 /tmp/gst_v4.py', timeout=20)
    raw = ''
    if isinstance(result, dict):
        raw = result.get('output', '').strip()
    cdp_state = None
    try:
        data = json.loads(raw)
        s = data.get('state', '')
        if s == 'true':
            cdp_state = True
        elif s == 'false':
            cdp_state = False
    except Exception:
        pass
    if cdp_state is not None:
        return {'checked': cdp_state}
    try:
        tree = env.controller.get_accessibility_tree()
        if tree is None:
            return {'checked': None, 'error': 'no_accessibility_tree'}
        tree_str = str(tree)
        label_lower = toggle_label.lower()
        lines = tree_str.split('\n')
        for i, line in enumerate(lines):
            if label_lower in line.lower():
                ctx = '\n'.join(lines[max(0, i - 3):min(len(lines), i + 4)])
                if re.search('\\bchecked\\b', ctx, re.IGNORECASE) and (not re.search('\\bunchecked\\b', ctx, re.IGNORECASE)):
                    return {'checked': True}
                if re.search('\\bunchecked\\b', ctx, re.IGNORECASE):
                    return {'checked': False}
        return {'checked': None, 'error': f'toggle_not_found: {toggle_label}'}
    except Exception as e:
        return {'checked': None, 'error': str(e)}

def get_bookmark_folder_and_dnt__33cb989983633f67656c84eb11865091_qw35sft2_9957cbaa(env, config: dict):
    """Get bookmark bar folder names and the Do Not Track setting."""
    result = {}
    os_type = env.vm_platform
    if os_type == 'Windows':
        bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Bookmarks'))")['output'].strip()
        pref_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Bookmarks'))")['output'].strip()
        pref_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif 'arm' in platform.machine():
        bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Bookmarks'))")['output'].strip()
        pref_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
    else:
        bm_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Bookmarks'))")['output'].strip()
        pref_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    bm_content = env.controller.get_file(bm_path)
    if bm_content:
        bm_data = json.loads(bm_content)
        bar_children = bm_data.get('roots', {}).get('bookmark_bar', {}).get('children', [])
        result['folder_names'] = [b['name'] for b in bar_children if b.get('type') == 'folder']
    else:
        result['folder_names'] = []
    pref_content = env.controller.get_file(pref_path)
    if pref_content:
        pref_data = json.loads(pref_content)
        result['do_not_track'] = bool(pref_data.get('enable_do_not_track', False))
    else:
        result['do_not_track'] = False
    return result

def get_chrome_appearance_and_dnt__197c5f21c248c0028a67e57d9193addd_qw35sft2_0f0a9cca(env, config: dict):
    """Get Chrome color scheme and Do Not Track setting from Preferences file."""
    prefs_path = '/home/user/.config/google-chrome/Default/Preferences'
    try:
        content = env.controller.get_file(prefs_path)
        if not content:
            return {'error': 'Preferences file not found', 'color_scheme': -1, 'do_not_track': False}
        data = json.loads(content)
        color_scheme = data.get('browser', {}).get('theme', {}).get('color_scheme', -1)
        do_not_track = bool(data.get('enable_do_not_track', False))
        return {'color_scheme': color_scheme, 'do_not_track': do_not_track}
    except Exception as e:
        return {'error': str(e), 'color_scheme': -1, 'do_not_track': False}

def get_chrome_search_dnt_state__67085c82c40b9ce759bdc0050139fed9_qw35sft2_10b9c364(env, config: dict):
    """Get Chrome default search engine name and Do Not Track setting."""
    import json
    import platform
    os_type = env.vm_platform
    if os_type == 'Linux':
        if 'arm' in platform.machine():
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    try:
        content = env.controller.get_file(prefs_path)
        data = json.loads(content)
        search_engine = data.get('default_search_provider_data', {}).get('template_url_data', {}).get('short_name', 'Google')
        dnt_enabled = bool(data.get('enable_do_not_track', False))
        return {'search_engine': search_engine, 'do_not_track': 'true' if dnt_enabled else 'false'}
    except Exception as e:
        return {'search_engine': 'Google', 'do_not_track': 'false', 'error': str(e)}

def get_chrome_lang_and_dnt__45fdf66a8d514412f8287b18b125d4f9_qw35sft2_14044afb(env, config: dict):
    """Get Chrome's interface language and Do Not Track setting."""
    os_type = env.vm_platform
    result = {}
    if os_type == 'Windows':
        local_state_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Local State'))")['output'].strip()
        prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        local_state_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
        prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            local_state_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            local_state_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
            prefs_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        return {'error': 'Unsupported OS'}
    try:
        local_state_content = env.controller.get_file(local_state_path)
        local_state_data = json.loads(local_state_content)
        result['language'] = local_state_data.get('intl', {}).get('app_locale', 'en-US')
    except Exception as e:
        result['language'] = 'en-US'
    try:
        prefs_content = env.controller.get_file(prefs_path)
        prefs_data = json.loads(prefs_content)
        dnt = prefs_data.get('enable_do_not_track', False)
        result['do_not_track'] = 'true' if dnt else 'false'
    except Exception as e:
        result['do_not_track'] = 'false'
    return result

def get_bookmarks_bar_urls__34222c255c19ca63c5f6efa3ee3ba731_qw35sft2_d29a6ea1(env, config: dict):
    """Get all URLs saved in the Chrome bookmarks bar."""
    import json
    result = env.controller.run_bash_script("cat /home/user/.config/google-chrome/Default/Bookmarks 2>/dev/null || echo '{}'", timeout=15)
    try:
        raw = result.get('output', '{}') if isinstance(result, dict) else str(result)
        data = json.loads(raw)
        bar_children = data.get('roots', {}).get('bookmark_bar', {}).get('children', [])
        urls = []
        stack = list(bar_children)
        while stack:
            item = stack.pop()
            if item.get('type') == 'url':
                urls.append(item.get('url', ''))
            elif item.get('type') == 'folder':
                stack.extend(item.get('children', []))
        return {'bar_urls': urls}
    except Exception as exc:
        return {'error': str(exc), 'bar_urls': []}

def get_chrome_delete_and_dnt__147c27e083ffc1ed15c90d12c99f1fef_qw35sft2_e3c218a8(env, config: dict):
    """Get combined Chrome state: auto-delete site data on close + Do Not Track."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\\\Chrome\\\\User Data\\\\Default\\\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    else:
        return {'error': 'Unsupported OS'}
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        cookies_val = data.get('profile', {}).get('default_content_setting_values', {}).get('cookies')
        data_delete_state = cookies_val == 4
        do_not_track = data.get('enable_do_not_track', False)
        return {'data_delete': 'true' if data_delete_state else 'false', 'do_not_track': 'true' if do_not_track else 'false'}
    except Exception as e:
        return {'error': str(e)}

def get_darktable_installed__bbd0b6cc65b6ed22ef194ee8993ef563_qw35sft2_022411c2(env, config: dict):
    """Check if darktable is installed on the system."""
    result = env.controller.run_bash_script("which darktable 2>/dev/null || echo ''", timeout=15)
    if isinstance(result, dict):
        stdout = result.get('output', result.get('stdout', '')).strip()
    else:
        stdout = str(result).strip()
    return {'darktable_path': stdout, 'installed': bool(stdout and 'darktable' in stdout)}

def get_docx_table_first_row__a4c465e58afc03255b4d0ddc46d5d525_qw35sft2_0ce40d24(env, config: dict):
    """Download Table_Of_Work_Effort_Instructions.docx, return table count, dims, and first row cell texts."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_rows': 0, 'last_cols': 0, 'first_row': []}
        last = tables[-1]
        rows = len(last.rows)
        cols = len(last.columns)
        first_row = [cell.text.strip() for cell in last.rows[0].cells] if rows > 0 else []
        return {'table_count': table_count, 'last_rows': rows, 'last_cols': cols, 'first_row': first_row}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_doc_tabstops__5d3adaf927fd3613be8530bc92e14de1_qw35sft2_e618f7eb(env, config: dict):
    """Read tab stops from all non-empty paragraphs in a docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/04 CHIN9505 EBook Purchasing info 2021 Jan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            stops = []
            for ts in para.paragraph_format.tab_stops:
                raw = str(ts.alignment).upper()
                align_name = raw.split('.')[-1] if '.' in raw else raw
                stops.append({'position_cm': round(ts.position.cm, 2), 'alignment': align_name})
            paragraphs.append({'text_start': para.text[:40], 'stops': stops})
        return {'paragraphs': paragraphs, 'count': len(paragraphs)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_italic_col__838d41f5c432af14f1eb478fa956e99e_qw35sft2_46821bef(env, config: dict):
    """Download Graphemes_Sound_Letter_Patterns.docx and check whether first-column cells in data rows are italic."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Graphemes_Sound_Letter_Patterns.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        if not tables:
            return {'table_count': 0, 'row_count': 0, 'data_rows_italic': False}
        table = tables[0]
        row_count = len(table.rows)
        italic_flags = []
        for row_idx in range(1, row_count):
            cell = table.rows[row_idx].cells[0]
            for para in cell.paragraphs:
                for run in para.runs:
                    italic_flags.append(bool(run.italic))
        data_rows_italic = len(italic_flags) > 0 and all(italic_flags)
        return {'table_count': len(tables), 'row_count': row_count, 'data_rows_italic': data_rows_italic}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_planet_table_rows__5135726bf1592f829eb43617d4e867fa_qw35sft2_f66e74b7(env, config: dict):
    """Get all row data from the first (planet comparison) table in the document."""
    import tempfile, os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/The Wonders of Our Solar System.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'rows': [], 'row_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return {'error': 'No tables found', 'rows': [], 'row_count': 0}
        table = doc.tables[0]
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        pluto_row = None
        for row in rows:
            if row and 'pluto' in row[0].lower():
                pluto_row = row
                break
        return {'row_count': len(rows), 'rows': rows, 'pluto_row': pluto_row, 'has_pluto': pluto_row is not None}
    except Exception as e:
        return {'error': str(e), 'rows': [], 'row_count': 0}
    finally:
        os.unlink(tmp_path)

def get_doc_tabstops__14c984754ad8a9510e10b62f00a4c316_qw35sft2_e918c2d3(env, config: dict):
    """Read tab stops from all non-empty paragraphs in a docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/04 CHIN9505 EBook Purchasing info 2021 Jan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            stops = []
            for ts in para.paragraph_format.tab_stops:
                raw = str(ts.alignment).upper()
                align_name = raw.split('.')[-1] if '.' in raw else raw
                stops.append({'position_cm': round(ts.position.cm, 2), 'alignment': align_name})
            paragraphs.append({'text_start': para.text[:40], 'stops': stops})
        return {'paragraphs': paragraphs, 'count': len(paragraphs)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_structure__638d8a63e2ebfdf518128094d93e23f6_qw35sft2_5c88a236(env, config: dict):
    """Download Graphemes_Sound_Letter_Patterns.docx and return table row/col counts and first cell text."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Graphemes_Sound_Letter_Patterns.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        if not tables:
            return {'table_count': 0, 'row_count': 0, 'col_count': 0, 'first_cell': ''}
        table = tables[0]
        row_count = len(table.rows)
        col_count = len(table.columns) if table.rows else 0
        first_cell = table.cell(0, 0).text.strip() if row_count > 0 and col_count > 0 else ''
        return {'table_count': len(tables), 'row_count': row_count, 'col_count': col_count, 'first_cell': first_cell}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_doc_tabstops__319c26794056ac5ca6a78fe65ba6a022_qw35sft2_e7ab72be(env, config: dict):
    """Read tab stops from all non-empty paragraphs in a docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/04 CHIN9505 EBook Purchasing info 2021 Jan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            stops = []
            for ts in para.paragraph_format.tab_stops:
                raw = str(ts.alignment).upper()
                align_name = raw.split('.')[-1] if '.' in raw else raw
                stops.append({'position_cm': round(ts.position.cm, 2), 'alignment': align_name})
            paragraphs.append({'text_start': para.text[:40], 'stops': stops})
        return {'paragraphs': paragraphs, 'count': len(paragraphs)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_header_align__83874f16b468982748e51b1480830a3a_qw35sft2_37bedf11(env, config: dict):
    """Download Graphemes_Sound_Letter_Patterns.docx and check alignment of header row cells."""
    import tempfile
    import os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    path = config.get('path', '/home/user/Desktop/Graphemes_Sound_Letter_Patterns.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        if not tables:
            return {'table_count': 0, 'row_count': 0, 'header_centered': False}
        table = tables[0]
        row_count = len(table.rows)
        header_centered = False
        if row_count > 0:
            align_flags = []
            for cell in table.rows[0].cells:
                for para in cell.paragraphs:
                    align = para.paragraph_format.alignment
                    align_flags.append(align == WD_ALIGN_PARAGRAPH.CENTER)
            header_centered = len(align_flags) > 0 and all(align_flags)
        return {'table_count': len(tables), 'row_count': row_count, 'header_centered': header_centered}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_dims__94b525d2310c47dd52401ee93f279d0a_qw35sft2_2214bd9c(env, config: dict):
    """Download Table_Of_Work_Effort_Instructions.docx and return table count and last table dimensions."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_rows': 0, 'last_cols': 0}
        last = tables[-1]
        return {'table_count': table_count, 'last_rows': len(last.rows), 'last_cols': len(last.columns)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_doc_tabstops__cf559754b29c29c8081fb54264f1ac10_qw35sft2_23c444b3(env, config: dict):
    """Read tab stops from all non-empty paragraphs in a docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/04 CHIN9505 EBook Purchasing info 2021 Jan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            stops = []
            for ts in para.paragraph_format.tab_stops:
                raw = str(ts.alignment).upper()
                align_name = raw.split('.')[-1] if '.' in raw else raw
                stops.append({'position_cm': round(ts.position.cm, 2), 'alignment': align_name})
            paragraphs.append({'text_start': para.text[:40], 'stops': stops})
        return {'paragraphs': paragraphs, 'count': len(paragraphs)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_bold_header__d0985f126b123d58872633f0350725ab_qw35sft2_3a1d6b53(env, config: dict):
    """Download Graphemes_Sound_Letter_Patterns.docx and check table existence and bold status of first row."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Graphemes_Sound_Letter_Patterns.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        if not tables:
            return {'table_count': 0, 'row_count': 0, 'header_bold': False}
        table = tables[0]
        row_count = len(table.rows)
        header_bold = False
        if row_count > 0:
            first_row_cells = table.rows[0].cells
            bold_flags = []
            for cell in first_row_cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        bold_flags.append(bool(run.bold))
                    if not para.runs:
                        rPr = para._p.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}rPr')
                        if rPr is not None:
                            b = rPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}b')
                            bold_flags.append(b is not None)
            header_bold = len(bold_flags) > 0 and all(bold_flags)
        return {'table_count': len(tables), 'row_count': row_count, 'header_bold': header_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_dims__f2fed6259df9a3cde848f19734b96eb2_qw35sft2_ac504bf7(env, config: dict):
    """Download Table_Of_Work_Effort_Instructions.docx and return table count and last table dimensions."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_rows': 0, 'last_cols': 0}
        last = tables[-1]
        return {'table_count': table_count, 'last_rows': len(last.rows), 'last_cols': len(last.columns)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_last_row__e0b81f05c0001516085ceb6f32f7b48c_qw35sft2_b0c3fdf7(env, config: dict):
    """Download Graphemes_Sound_Letter_Patterns.docx and return table row count and last row first-cell text."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Graphemes_Sound_Letter_Patterns.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        if not tables:
            return {'table_count': 0, 'row_count': 0, 'last_row_cell0': ''}
        table = tables[0]
        row_count = len(table.rows)
        last_row_cell0 = ''
        if row_count > 0:
            last_row_cell0 = table.rows[row_count - 1].cells[0].text.strip()
        return {'table_count': len(tables), 'row_count': row_count, 'last_row_cell0': last_row_cell0}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_timetable_remove_tutorial6__c0f66209d7bea00756990ac3a15deab7_qw35sft2_f9a85fc1(env, config: dict):
    """Read D5 (Wed 12PM) and G12 (Sat 20PM Tutorial 6) from Course Timetable.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/Desktop/Course Timetable.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        d5 = ws['D5'].value
        g12 = ws['G12'].value
        return {'d5_value': str(d5) if d5 is not None else None, 'g12_value': str(g12) if g12 is not None else None}
    finally:
        os.unlink(tmp_path)

def get_chrome_tabs_and_bookmarks__ba148d88a9035c4354cf728e6417abab_qw35sft2_bb9b51d7(env, config: dict):
    """Get both open tabs and bookmarks from Chrome for combined evaluation."""
    result = {'tabs': [], 'bookmarks': {}}
    try:
        tabs_result = env.controller.run_bash_script('curl -s http://localhost:9222/json 2>/dev/null', timeout=15)
        if isinstance(tabs_result, dict) and tabs_result.get('output'):
            tabs_data = json.loads(tabs_result['output'])
            result['tabs'] = [{'url': t.get('url', ''), 'title': t.get('title', '')} for t in tabs_data if isinstance(t, dict) and t.get('type') == 'page']
    except Exception:
        pass
    try:
        bm_result = env.controller.run_bash_script('cat /home/user/.config/google-chrome/Default/Bookmarks 2>/dev/null', timeout=10)
        if isinstance(bm_result, dict) and bm_result.get('output'):
            result['bookmarks'] = json.loads(bm_result['output'])
    except Exception:
        pass
    return result

def get_timetable_lec_color__462e6cf4d0037b83d0687915bbffba64_qw35sft2_157dbc7f(env, config: dict):
    """Read D5 value and fill color from Course Timetable.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/Desktop/Course Timetable.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        d5 = ws['D5']
        fill_color = None
        if d5.fill and d5.fill.fgColor:
            try:
                fill_color = d5.fill.fgColor.rgb
            except Exception:
                fill_color = None
        return {'d5_value': str(d5.value) if d5.value is not None else None, 'd5_fill_color': fill_color}
    finally:
        os.unlink(tmp_path)

def get_extension_project_state__ce07ffbd7d3847464df6caec281dec8d_qw35sft2_770a4b6d(env, config: dict):
    """Check the state of the happy-extension project in ~/Projects."""
    base_dir = config.get('base_dir', '/home/user/Projects/happy-extension')
    checks = {'dir': f'[ -d "{base_dir}" ] && echo "YES" || echo "NO"', 'manifest': f'[ -f "{base_dir}/manifest.json" ] && echo "YES" || echo "NO"', 'background': f'[ -f "{base_dir}/background_script.js" ] && echo "YES" || echo "NO"'}
    state = {}
    for key, cmd in checks.items():
        result = env.controller.run_bash_script(cmd, timeout=15)
        if isinstance(result, dict):
            output = result.get('output', '') or result.get('stdout', '') or str(result)
        else:
            output = str(result)
        state[key] = 'YES' in output.strip()
    return state
