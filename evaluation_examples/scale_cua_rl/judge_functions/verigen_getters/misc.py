"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_docx_para_alignment__099018697289346087fe050b42a1dd16', 'get_compose_fields__457e68f7ef86811b4470193bf0ca1e7e', 'get_docx_first_line__2308ac81a92b82ba916e71c364d68457', 'get_docx_default_font__450e6c1e4a27707dc54b625806c54889', 'get_docx_heading2_alignment__131e78d644cdd4aebe878331238d0634', 'get_writer_font_sizes__5b4d6572a569badb37aa2fd25bab4b6e', 'get_docx_footer_page_numbers__9de26a1b58798be1fa2f4b6b3d9e9ec9', 'get_docx_styled_merge__3faeea9d9342eb480cdf899557de9715', 'get_docx_first_para_alignment__f845858287fa4f02431e180c3025ad53', 'get_docx_reverse_merge__08f1db516ca13465963b37a9fdcaf66c', 'get_docx_selective_merge__1d49dccb31c44fca7886c9c825e7148f', 'get_docx_para_alignment__4ed04700852dd5a18fbdc223aac7b1c3', 'get_docx_text__896d76d765c6621bc8feff2ba14e6be9', 'get_power_settings_state__e637b95c625a88b3bb2cee1c77022b86', 'get_docx_title_alignment__0a4a5cad7031961d3e6b497ec2c282f4', 'get_check_utils_py__597b85b971b86604637b8eea77967812', 'get_docx_font_info__2042c7daa4f47bea2c4fb2c5e4661307', 'get_docx_text_formatting__92a351c44e937610849c0af1deb80b89', 'get_header_italic__ee9130b48c4f13e968a7877f09e8d206', 'get_docx_italic_color__575ab439795c172e41189face7efb98a', 'get_docx_first_para_font_size__0625ebacf8174a330d7b47b181788578', 'get_background_color_state__0fc932ce1ad9456b6d290262a9551736', 'get_docx_para_font__ec9fd979c9f4ede06619498c1e9b1202', 'get_row_hidden_state__18d397a77a413967bac86e92a1394add', 'get_docx_all_fonts__e754543a404f008975bc7bab115cbf09', 'get_tally_rows__89c74c0ac9693815f604c967708f8f82', 'get_docx_station_lines__dc2cb3918367317ad49987158bb0de7d', 'get_check_first_line__35d6e56c2b031e6f068f1907dd34bebf', 'get_docx_footer_page_numbers__806bfdf7c0ee49dcce08f91a2edad9e6', 'get_docx_text__c53695cdc8b25b507dc6019962dcaf31', 'get_writer_highlight_check__6b6e5048ebb3170359c47d1ce863bb95', 'get_vtt_conversion_state__9ac6e58c407c5a10e99c8a824933ab2f', 'get_docx_page_numbers__18ec7a61b2f5504dc6e64dce31cd1e00', 'get_user_password_home__6fbd0fb5ef1077b146e8ceece522a6b4', 'get_writer_first_para_bold__974838092a7939215ca9c77d0b32d057', 'get_docx_heading2_fonts__3ee05560ea0a8eae187c5c1b3f30c5d3', 'get_docx_title_text__ca26e505bd69174991058c87e74f2a98', 'get_docx_answer_line__12b6d7e1507f0e7212c8d9d4764a8a2b', 'get_docx_italic_bold__7532d291294fb4f3bb8e25015109f986', 'get_docx_title_format__c874bb9d2ad63d6767d6a39e179946a5', 'get_docx_first_para_font__53ba1671387b371fc528911be537db49', 'get_notification_dual_settings__0a0313fec36cad178ad94d7913e78201', 'get_writer_font_check__2a303ca90381d050b859a1a253c0648a', 'get_writer_title_italic__088f05c0fd2c52a92bc3dc8f8cc6719c', 'get_docx_footer_page_numbers__3205af8e0dea88590c54f355855d7cbf', 'get_contrast_comparison__a61095a992c03ac5e17b282d64e6f2d0', 'get_docx_italic_font_name__f2fb15859784992fd6cf15bd171d4b91', 'get_settings_speed_blocksize__4b546ac4b7d86d470dd227a5c08930e6', 'get_docx_first_para_alignment__ae9bdf892a1bf6b8ab5ef93251d291e0', 'get_docx_para_bold__7a1ffa51886994c5f3667d404d37d94d', 'get_docx_all_bold__aef5992d33ef93f956afe1a194ef0bce', 'get_booking_form_state__cf1c94ef66b2087660e9bbd9fc7c9d7e', 'get_docx_title_font__1b8a8292a0fd2eb13bf16e945356170a', 'get_docx_title_alignment__4f408a61496e9bff25fa5074f0f2e380', 'get_compose_recipients__9a5e3321e7daf20a2f35cb35407bc19e', 'get_writer_default_font__0f58b322b0ad4bcc09c937490ff9bfeb', 'get_docx_all_fonts__f0d9875ec4c08351c2fb33918a5793ea', 'get_main_caption__909b46690f7fc746deb97149427ce266', 'get_docx_first_para_alignment__58a0796c274d2d883482eaaa895cfb2f', 'get_writer_title_fontsize__22037d9b737ef0f8bf33646ac10692bd', 'get_docx_title_alignment__6fdb2c432a1af94ddff3a56f751c6296', 'get_budget_page_state__b76001106213215f01706a09eab507db', 'get_check_uncomment__cd3236440d08e0a8882e8811427d01d3', 'get_docx_delimiter_check__0a313c15fc4fb10c48c753cea47e4527', 'get_docx_font_check__c5ed5054f5dbce73ddb0483f71ff6f1a', 'get_docx_all_text_case__0c637a2f0d91975120dc8e0a76554911', 'get_docx_heading_alignment__e69f89d4207a469ba2d4f747df580e7b', 'get_docx_page_orientation__7f5136b7f1dcbd0181f2d126bcac0caf', 'get_settings_blue_color__be59b6debe1fbee8fc53696bdb2273b8', 'get_docx_italic_status__bf5294aa1b721bbae785292729cb8020', 'get_color_and_size_state__043ea3777ab786b45abdbd440171ffc3', 'get_docx_first_para_italic__d7547f75d107bc34f4a878c280fffbc1', 'get_docx_heading_alignment__de2b396948b313eaa0067057e04e4e77', 'get_docx_title_formatting__6146cd6c34c0bc769a02e73ec8ef6f70', 'get_docx_hint_line_spacing__7fbc318f6c5b72e9f18e5cf28491b9c0', 'get_docx_footer_page_numbers__ad71c47cc02b102c764d5154ce0d73c0', 'get_docx_all_fonts__2ae82910da68002280a1510e1ca61e99', 'get_docx_text_replace_check__3df1073706339202301b753c7942ef7f', 'get_writer_title_alignment__3a67dac16aed3d74bffb9cb95c614713', 'get_rows_hidden_state__cb276526eec2ec94231e7dedd20b5069', 'get_docx_title_bold__94187b1698fde7f55883e20a5cb9822f', 'get_venv_check__3936cc1f2257a01f73d410787af3a332', 'get_budget_reservation_info__a4a15f9cf85bf343423f5b15bc460e50', 'get_settings_fps__6d1f5614df32aa00fba65a5a2ca43cb2', 'get_docx_footer_page_numbers__b8220411d07c53e298b6db0365205e12', 'get_docx_default_font__fc8e11a362be23f66acd12dd6bff4956', 'get_docx_heading_bold__0d3391a68f9d89ec2a40039cb5c06d2d', 'get_docx_answer_line__0636655694efccc23299c2ca2f236f1f', 'get_docx_text__72fc0161fa321ae4bf01e0be489f2375', 'get_docx_text_case__0a695088802a9cb1473c98bdb4ec1b4f', 'get_docx_first_line__2130fa798c3c78ff755136f88d8b3efd', 'get_docx_default_font__536a5616674c51027b6cacf28aa7ecd9', 'get_delta_miles_checkbox__f9b3359b12066a252ed8849a9876d835_qw35sft2_b6e19eb1', 'get_delta_form_state__45facd29054b2deec3804036a06322c9_qw35sft2_dd814e40', 'get_delta_form_state__607589a24fcd7d8000928f466f0dd577_qw35sft2_503fe9b7', 'get_delta_destination_field__696c0a5477f795f89ba4614ef9b530f9_qw35sft2_06ddfc62', 'get_history_and_dnt__43a7a4d8fb3eedcedaf50f5a89a2ea93_qw35sft2_3461e6ae', 'get_font_and_dnt__baedb9f4aec1e7501564ae9f37c335ac_qw35sft2_30331902', 'get_frame_at_3s__81411ca26f8f6e6edae115d08d9ab086_qw35sft2_cd74764a', 'get_freeze_panes_state__10368a4827622b87edd2c56b2f249e0f_qw35sft2_2e501aba', 'get_salesrep_jan_total__5238964d8e657c42c8a059231010b871_qw35sft2_67ceaa27', 'get_income_gross_with_total__9bcff5517fb7493e61233c0a423569f9_qw35sft2_58941de0', 'get_employee_ages_avg__151640eb3d43a7eb22c551550e103953_qw35sft2_74381784', 'get_vlookup_f2_single__dba068db6ac4d383c892aae7e4d7fc62_qw35sft2_44448389', 'get_total_label_and_jan__4e23e64ce56ffa608b1e8fd2866f5ae0_qw35sft2_58708e5a', 'get_header_bold_format__988aa3731e056dbe2f61a28b2135e9ae_qw35sft2_7afefb80', 'get_monthly_totals_row__aff1e90d6f0bac6579c8d1279e6c6c3e_qw35sft2_721b5769', 'get_weekly_sales_profit_total_row__50f7148385ea6d504119cd2e40f7d1be_qw35sft2_e95480fb', 'get_employee_split_sorted__5d1c2e4438b4aac3fb5db710d29300f5_qw35sft2_bcc33200', 'get_period_rate_max_in_d1__14298c6976685b9e8d10b60e8490521a_qw35sft2_14e7047e', 'get_maturity_sorted__6a0b6ecdb454d8389c50fe9e251b1580_qw35sft2_e71a2e21', 'get_salesrep_label__d6d9e5e1d2ca3d967d1969fd1a1e0c2c_qw35sft2_9789eb01', 'get_seqno_and_total_row__f9ba38ec0e00ec28c9c0ef5057c79916_qw35sft2_c93bd6f3', 'get_employee_ages__add589b113a1be1eec737c1ab22661fd_qw35sft2_b15815f6', 'get_vlookup_f2_f4_rows__22ff441442432fb67bd17ef3eb97a292_qw35sft2_21d8f812', 'get_income_net_sales_gross__c6e3aa9fd1cbfd34b4473b9e8f31e349_qw35sft2_7e167ef8', 'get_total_row_state__5ac69261c8fb294b16f18049ece06ce9_qw35sft2_313652f6', 'get_row_hidden_state__93bd58adf830a704a8880938e35d28c4_qw35sft2_83a00ea9', 'get_weekly_sales_sorted_by_profit__99a00a677b41fb8db78ec79d5381f0b2_qw35sft2_eae34b45', 'get_employee_split_bold_header__fc565b338fbb3db18703aef2eaa623cb_qw35sft2_4d4d17f1', 'get_ramp_accel_diff__e64652ba3bf9a74b525e231bf3f391d2_qw35sft2_6edf4814', 'get_period_rate_sum_c26__fd409a83b18f09cf45613c5173377d4d_qw35sft2_ec0397ab', 'get_maturity_total__461252d0c8044977d30ecc240f2dc9cd_qw35sft2_08c41667', 'get_freeze_panes_state__ae87df329e6b7ba3f255d9a2acdfeda6_qw35sft2_8fbda1a7', 'get_seqno_and_sales_sum__00713006369c01f94f764c7eb6008543_qw35sft2_bdbb170f', 'get_salesrep_last3_totals__ce77d41e3a7b3efa48e4655c22c8aa27_qw35sft2_dac8c6be', 'get_vlookup_f2_f12_all__202132c1158925d79d1ba222174d8f66_qw35sft2_e41cb20f', 'get_income_net_sales_and_cost__43eee44b990a73c4b5a6e2bc138833a7_qw35sft2_975e4d4c', 'get_strikethrough_and_transition__b374fc122f888120a74e40bee1da5133_qw35sft2_e3f3d361', 'get_title_font_props__d1408966fdb4fd77bea2fd3e21b1ee11_qw35sft2_60a6ab82', 'get_strikethrough_and_italic__be4183a9c6982ba683532f14bc50d4bc_qw35sft2_3f4cf417', 'get_picture_heights__eff65ebb8ed6d102db81f31da1a823ed_qw35sft2_c0347b22', 'get_title_font_props__f083039072ffd98eec2d101bb28501a7_qw35sft2_a86e06fe', 'get_strikethrough_state__7f3881fcd58bcba327d08f20ed190dcc_qw35sft2_a43cff2a', 'get_writer_footer_body_spacing__e86d756354df667eec5758969f058362_qw35sft2_2e25e64c', 'get_title_run_font__ac8fd7a53908049060d14f1a620d5f79_qw35sft2_4c453cf5', 'get_para0_underline__0ee57f85e0ff64364a5561938bc94f89_qw35sft2_eeee8dff', 'get_docx_last_line__9ae69d0ef002958997c53c817727d757_qw35sft2_1cd0991b', 'get_docx_subscript_and_title_italic__5faa01d8a062ee366a1e56b39b81f8d9_qw35sft2_8cf73d1b', 'get_writer_basic_fonts__cb1a08b69b2e8e76cab6534a166b3ea6_qw35sft2_3329d88b', 'get_docx_lower_and_page_nums__4155f473baa624a7c08f0169367bee57_qw35sft2_721e17f7', 'get_writer_titlecase_center__c181651f58b43b69af608583a0523891_qw35sft2_1b32c475', 'get_writer_font_align__2b822025bba1f4240036b27d0339b871_qw35sft2_ac7f2755', 'get_writer_heading_body_align__d846f8a9d2cb6584815a6f17cfb8e80a_qw35sft2_0cdf8986', 'get_docx_first_para_strike__efe155f44671b1c09a26d5fefbd2fc44_qw35sft2_f933e3da', 'get_docx_italic_font_size__b8546f99d1a88f06fb0b6ffbfa13cf55_qw35sft2_e5a0021c', 'get_word_font_colors__336f1c4245e680e1aace133fb243059a_qw35sft2_3479f659', 'get_odt_highlight_italic__e4e69693f939d9cc32e39d11fb21f92a_qw35sft2_45bc9ddb', 'get_docx_three_para_spacing__e4c36138929fcbc781e762a8996148c1_qw35sft2_3d322a41', 'get_writer_extended_state__d40d48ebd32c5cdd693a8c4ab4565d6c_qw35sft2_1970bdc2', 'get_docx_spacing_arial__d472f99c5ab0bbd719c99178119f321f_qw35sft2_85781c1c', 'get_writer_break_and_notes__4e9ea3f60c13559abc5abe689f4b9ea2_qw35sft2_57942942', 'get_writer_default_and_list_fonts__12134317917b6b593b5731418bb79a41_qw35sft2_75bbc093', 'get_docx_first_line__d29d0e1e9fda0acc872853f3a63d2906_qw35sft2_69a56e4c', 'get_writer_footer_title_align__54873b40ed0174b0628f230c8ad47868_qw35sft2_ecff64b7', 'get_docx_lower_and_bold_title__5844d45da47ec52046a23219bd5db770_qw35sft2_ab4c298e', 'get_docx_subscript_and_title_center__49f57bac6148724412fcfb3eb162539d_qw35sft2_4fe9563d', 'get_writer_titlecase_italic__7a4ad09a1377c8f0eb798dbce87a7dce_qw35sft2_1af9ffcd', 'get_docx_multi_para_strike__d8440cb16db3c65fff5a2b530ef73072_qw35sft2_8f67219f', 'get_docx_italic_size_underline__dce073995a5927fc743181d7c02c0659_qw35sft2_3d981ede', 'get_docx_line_spacing__c985596bc3953d61c7d274fc30ba4960_qw35sft2_14feb6d7', 'get_writer_font_size__a567f4bdee19e51d5ed6a587d8d38f9d_qw35sft2_ae1290b0', 'get_writer_heading_align_size__a3c76448cfe24f131aa7bba8054b5d03_qw35sft2_7bb7f6c5', 'get_word_font_colors__d89e455c3ca4dbc17e773411cf8f66df_qw35sft2_0501c7f4', 'get_odt_highlight_font__582984f0118e076c9ea08b5f36393d3b_qw35sft2_afca0b7a', 'get_writer_three_state__c8e0b7dd7367f091fd322e812e9986d0_qw35sft2_98c08613', 'get_docx_spacing_fontsize__122cebb30dd4bbe8d6067021381568cf_qw35sft2_355b80cf', 'get_writer_break_and_font__52445c32d012b75894f7e753d9bb73ed_qw35sft2_9ef5f0c7', 'get_title_alignment__6435e9b69d0bbcb863c89964e3084688_qw35sft2_568a8a6b', 'get_writer_titlecase_underline__d22455594f3cb66658bea4571465f8f1_qw35sft2_d8cd9c76', 'get_writer_footer_pageno__035b6f0ae922d6e2d1f24326ec60c904_qw35sft2_865faa47', 'get_docx_text_all_upper__cbc878ff0e7a736898047bffbd6f85f6_qw35sft2_5ca44771', 'get_para0_bold__42cf2947548fcbcc72e5e44a50eb60dd_qw35sft2_99cf7f57', 'get_docx_train_records__aa11b4f1fde2d6cc216fd8dac61371d5_qw35sft2_1468b0d6', 'get_writer_font_and_alignment__70647269807c892c9cd9454bb994336f_qw35sft2_5d687801', 'get_docx_subscript_and_bold_heading__fa57f784af7171fa0e677f21369b4177_qw35sft2_c13dc097', 'get_docx_italic_size_16__e2de398ff1a46d21a89c93f18cb16653_qw35sft2_5af2a69c', 'get_writer_font_bold__3a9eb9ceb78a9dd18ebdb09636bdbe29_qw35sft2_786d155a', 'get_writer_heading_align_font__fb136ac47c87dcc7e787c09564b245bc_qw35sft2_7c722193', 'get_docx_strike_bold_chain__30608d866c349bd16efe53a1c72a27d7_qw35sft2_95f689f7', 'get_word_font_colors__e97799f0add4268fcc31cdd1e95ac277_qw35sft2_a5782795', 'get_docx_mixed_spacing__ffe17a8bd50575f09eec01a68fed60c6_qw35sft2_6f5cace1', 'get_odt_highlight_strikethrough__a30ceeb68a8eeb0b8a512ab61a2387e1_qw35sft2_8ca7650f', 'get_writer_ref_state__dab123b7868d167252ed777095add0ba_qw35sft2_03ce8375', 'get_docx_spacing_bold_intro__7ab11a38694adeb4332550a69fd5abfe_qw35sft2_57061efb', 'get_writer_break_and_title_italic__0f7fdf2d11eca0ddff6ddca8ba1c7a92_qw35sft2_9f93c670', 'get_docx_sentence_case__121c5b95c7fa376e725e2a74f1f18a20_qw35sft2_0ffe6d01', 'get_writer_footer_title_italic__6a6cc374d796e945ca1edd575879589a_qw35sft2_e128fbf1', 'get_docx_last_line__8efbbd8cc0ef91addd7e2f867e2ff2b5_qw35sft2_cefd8c72', 'get_docx_last_line__e6176eae3a47bbdd8f12a06b8b082ed4_qw35sft2_cb11f3e2', 'get_writer_font_and_pagebreak__e2b27ec72c919c5a4bd4f151e7ff64b6_qw35sft2_df0dbf5e', 'get_writer_titlecase_doublespace__6e2b5ade87cb088089f67dc5070eb77b_qw35sft2_362038d8', 'get_docx_subscript_title_and_body__cb7672ac12071a9e4c398d11dc4c470e_qw35sft2_96ea8482', 'get_writer_font_italic__c3632663016de6f20637b59e83145d05_qw35sft2_4378d365', 'get_doc_page_breaks__b7c367c074a6362d7d9c85a08867a367_qw35sft2_dbc4af6a', 'get_docx_italic_and_title_size__792e97b1e8b831bfb11064f6cde55b00_qw35sft2_660a7e52', 'get_writer_heading_align_italic__edb463affe3b938af85226bb52cb1b03_qw35sft2_b85ec2f9', 'get_para0_italic__d2318c984412daccf01ad6c479224170_qw35sft2_1c924744', 'get_docx_last_para_bold__c977b8be915d8296c568265b630cf188_qw35sft2_851b1d39', 'get_word_font_colors__3d9f327a66b9025c03e34cafe8c88fe3_qw35sft2_ef29a4bb', 'get_docx_spacing_italic_conclusion__bb617d210691b208d4c61100467118c1_qw35sft2_e414146c', 'get_writer_ref_footer_state__78615d30a3700f795271d9e8e63e0686_qw35sft2_1ff19677', 'get_writer_page_break_count__e042d7b442613a05635c554b051c41b1_qw35sft2_7b22bea1', 'get_odt_highlight_fontsize__4ce1d37833ab4f6aa409d7454a343799_qw35sft2_cedb6e29', 'get_doc_line_spacing__80aefd2088f27d1be89724efb49ed091_qw35sft2_0db69134', 'get_writer_footer_and_header__5d6628525b72a5805637db6ac940e950_qw35sft2_bc1fbc77', 'get_docx_subscript_and_heading_underline__b75f28775a1b3fe6faeb633d14a05fab_qw35sft2_913900f0', 'get_writer_titlecase_bold__2777e85b511814778d7406b21395647f_qw35sft2_643ea196', 'get_footer_page_numbers__699ab5651848f548e06466de9777d875_qw35sft2_dbe640d9', 'get_docx_dedup_state__f922eeeed3d49013fd1e13103dbfb120_qw35sft2_7be396a3', 'get_writer_font_and_footer__d1ae06a2cb3f08c6662cf614ce58e285_qw35sft2_25efcd2d', 'get_docx_italic_bold_size__57bf430ce2e461a40bf234942a8ba4ad_qw35sft2_fe2c4902', 'get_writer_font_size__5934b0de41c172c0d1662adc3bbc874f_qw35sft2_a3711a32', 'get_writer_heading_align_body_bold__e1c39faa270c25698b597d87598bba49_qw35sft2_7ed4bddf', 'get_para0_font_size__7211ff5a62d6ff910c67632b31846d1c_qw35sft2_982eaf9d', 'get_docx_second_para_italic__aa396abb2d9562e18e97067a7c8c6fe9_qw35sft2_ddda8fe7', 'get_docx_spacing_center_intro__5f00d786b29854e31074fa3cfeefea1b_qw35sft2_3e7fe962', 'get_writer_citation14_state__21b44c209dfc410a39a6375bc0042826_qw35sft2_983282d0', 'get_invoiceGES_in_problematic__cea86e2d544f9d987b2adcf034f36c4a_qw35sft2_7a0e7baa', 'get_docx_with_header__2354b786c1a1e08afc94d3d722cbb7a6_qw35sft2_bdd3f828', 'get_main_py_first_line__7b62591361b18caf2ed6aa916d0c9cf5_qw35sft2_61a5d212', 'get_docx_text__ed93db9128130b7a782216ac4507a0ca_qw35sft2_c07f733b', 'get_docx_first_para__40a84e6a55059151959a2b459d4c099e_qw35sft2_89c54b9d', 'get_sar_disk_report_state__fc8ac0599dbb02ca56427b3ae265c796_qw35sft2_6bd8bd41', 'get_docx_writer_state__8ae851f2fbec3e62c3a4b8c28bfb6c8b_qw35sft2_a3812e2a', 'get_settings_snake_size__85e39531da2c848a1afaf10e25ba3dad_qw35sft2_f06e060f', 'get_invoiceTII_in_problematic__7ecb63a4609903641df39f4754b7248f_qw35sft2_84d08576', 'get_book_copy_in_documents__55d3e622f74c4c6d9051e6e358a64ecd_qw35sft2_a3d090d3', 'get_combined_state__3001e2e091ccfaef8fccec23c5f15501_qw35sft2_9a14dca8', 'get_tally_book_amount_sum__289ae9cd25fb72b631d3d5c97c4b9f82_qw35sft2_fb6205ee', 'get_paper03_and_year__d508d903f596b0ecb03360bdda6892d4_qw35sft2_17558f40', 'get_docx_duration_line__48db15ca4dac03d9c5bae1d76e883852_qw35sft2_7725eb8d', 'get_desktop_listing__00bb03d81e57baa40bd1c86ce0b1574d_qw35sft2_5b6e050b', 'get_sar_cpu_report_state__5dfc721833abb9e182bb18218c19d630_qw35sft2_32121cc6', 'get_docx_gemini_paragraphs__abfc7551de3c0f45f173130d54033329_qw35sft2_a085b892', 'get_docx_writer_state__5d4910d5e528b018db8af7469dad985a_qw35sft2_4b976908', 'get_settings_fps__751d0a9ef75d8ec2c4c27e57248ad37c_qw35sft2_74691e56', 'get_pandoc_install_status__79ba7b4d76657204de8c0df9dc21c02e_qw35sft2_452b3542', 'get_invoice243729_in_problematic__da5edbd81a2d596fb7eb6ef7b52c2151_qw35sft2_e46fc3b3', 'get_conda_install_state__d5abdc9092a51546a82222c972d5edf6_qw35sft2_98581235', 'get_paper01_and_count__101069bcaae5b37861b743495ae4859d_qw35sft2_6d9a0ea7', 'get_desktop_listing__b0e8ccdfade877a9b91932809683825c_qw35sft2_bfc0cbaa', 'get_tally_book_last_row__97f6f4530ffb08062fee8d0568f59984_qw35sft2_3aa47332', 'get_copy_move_state__ec373221258728f3163857f2bdfd353a_qw35sft2_d2ffe179', 'get_clock_dual_settings__443802668a0bfacbf0757782ecb2ad43_qw35sft2_39fa83f8', 'get_accessibility_large_text_no_animations__855f455c8e49beda6b556868bdd9b753_qw35sft2_c390975a', 'get_gnome_favorites__41851786639af4a56b165470f26fecab_qw35sft2_bbbb11b6', 'get_volume_level__d6b3f95410b1cde5a46558c01aaafc46_qw35sft2_5eb74f22', 'get_timezone_and_clock__c55229b53381e6587f722d00658cbd02_qw35sft2_79b14818', 'get_rename_and_sibling__9cc29f7146497c596e333ae0ebaf5848_qw35sft2_82e8f5b4', 'get_notif_clock__accf3abf7a1d21f7f68eefd98414a494_qw35sft2_2b863f6d', 'get_power_dim_and_blank__4a865a2b73390cd1be0bc5f929d03ff5_qw35sft2_fd6f3625', 'get_php_stats__a63cdece1fd90a60a3e737af233f1764_qw35sft2_2b1cee42', 'get_move_failed_notebooks__96695aa3604bded58ed2ec7f3c72c8de_qw35sft2_932f1c29', 'get_output_and_backup__9fff5b0d3da3f15287b691b4b1944f6b_qw35sft2_e025a80d', 'get_rename_and_move__eb200f79b699eb5d233b5379c7e1ea0e_qw35sft2_f80e9fe2', 'get_volume_level__2ac92884cbfeefbd5a2d6929c20839d7_qw35sft2_ff2365b1', 'get_restore_copy__f767d224a0cdea04b11256d9cffabd41_qw35sft2_5940f6b1', 'get_copy_and_count_fails__8f62d3b2706e8f98a78e22aea189146a_qw35sft2_c7a85b79', 'get_timezone_and_ntp__397a38a7528d157681d1e57316ff4073_qw35sft2_6e2cf5df', 'get_volume_settings_state__1201dde5af3c06b515ed16eae4fad685_qw35sft2_6b004e3e', 'get_gnome_favorites__f6775b09b9b85d12e5fafc15ee8143de_qw35sft2_cc93bd8a', 'get_power_triple_state__570bc607f27286d53f7b1a8eaf2af624_qw35sft2_a3f7abdb', 'get_notif_sounds__139fa3103c161a5284bcac4b5dc4f632_qw35sft2_a18ce23b', 'get_gnome_favorites__fafd25459d3a377f8d0c6d91b93474de_qw35sft2_5897609a', 'get_power_dim_and_suspend__6cd1c7feeea3fc52512c1552e8d5b3b2_qw35sft2_dddf5aff', 'get_copy_and_rename__44eff79b79182e6eeb0c7debb1bc9f8c_qw35sft2_de77ce1d', 'get_accessibility_two_toggles__8ee8727a5fa2032ce96fa89a79ac9fab_qw35sft2_2815a142', 'get_eml_backup_state__d81f9d153146f886b82c4131d09d1ccb_qw35sft2_e1e59461', 'get_bills_flag_and_tags__ed9097e7c75920408b4536024b8da2a4_qw35sft2_80b3d527', 'get_bills_full_state__d21b0cc63afcca51c02fdd16b5d862e0_qw35sft2_0be511a0', 'get_eml_count__2731b9abd5cfbad9ed4df8aae737addc_qw35sft2_d185deda', 'get_bills_flag_state__8e8eeb1588f1109e98ccb02a0faa787c_qw35sft2_87acafa7', 'get_eml_listing__dfeb48225188ee18fb4de9d6f0048829_qw35sft2_71b752de', 'get_bills_and_filter__44c5673d772aeaf1fd4e6dcc8a32111d_qw35sft2_a2ee062f', 'get_ext_and_multi_settings__afc2fb3b68b53df8d476e05f07933aff_qw35sft2_71352dcc', 'get_ext_and_fontsize__fe14a817663aeade8a20cb0f5baad2b6_qw35sft2_e0fd936f', 'get_ext_and_wordwrap__324a83eafb9ff6ed07fafe7e199af4d5_qw35sft2_1f593828', 'get_ext_and_settings__bb23bc65a82d76a52cba97bb0a6bf9bc_qw35sft2_923dbd2e', 'get_dual_ext__95ee9f441436911870f312520ed4f195_qw35sft2_08027149', 'get_ext_and_wordwrap__6ac7db98d9670abbcf37d963dc27bc84_qw35sft2_f286792b']

def get_docx_para_alignment__099018697289346087fe050b42a1dd16(env, config: dict):
    """Get alignment of each non-empty paragraph in a docx file."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        alignments = []
        for p in doc.paragraphs:
            if p.text.strip():
                align = p.paragraph_format.alignment
                if align is None:
                    alignments.append('left')
                elif align == WD_ALIGN_PARAGRAPH.CENTER:
                    alignments.append('center')
                elif align == WD_ALIGN_PARAGRAPH.RIGHT:
                    alignments.append('right')
                elif align == WD_ALIGN_PARAGRAPH.JUSTIFY:
                    alignments.append('justify')
                else:
                    alignments.append('left')
        return {'alignments': alignments}
    finally:
        os.unlink(tmp_path)

def get_compose_fields__457e68f7ef86811b4470193bf0ca1e7e(env, config: dict):
    """Get compose window fields from accessibility tree."""
    try:
        tree = env.controller.get_accessibility_tree()
        return {'tree_content': tree if isinstance(tree, str) else str(tree)}
    except Exception as e:
        return {'error': str(e)}

def get_docx_first_line__2308ac81a92b82ba916e71c364d68457(env, config: dict):
    """Get the first non-empty paragraph text from a docx file."""
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            return {'first_paragraph': paragraphs[0] if paragraphs else '', 'all_paragraphs': paragraphs}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_default_font__450e6c1e4a27707dc54b625806c54889(env, config: dict):
    """Get the default font and most common body font from the docx."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        default_font = None
        normal_style = doc.styles['Normal']
        if normal_style.font.name:
            default_font = normal_style.font.name
        font_counts = {}
        for p in doc.paragraphs:
            for run in p.runs:
                fname = run.font.name
                if fname:
                    font_counts[fname] = font_counts.get(fname, 0) + len(run.text)
        most_common = max(font_counts, key=font_counts.get) if font_counts else None
        return {'default_font': default_font, 'most_common_font': most_common}
    finally:
        os.unlink(tmp_path)

def get_docx_heading2_alignment__131e78d644cdd4aebe878331238d0634(env, config: dict):
    """Get alignment of all Heading 2 paragraphs in the docx file."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        heading_alignments = []
        for p in doc.paragraphs:
            if p.style and p.style.name == 'Heading 2':
                alignment_str = 'UNKNOWN'
                if p.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    alignment_str = 'CENTER'
                elif p.alignment == WD_ALIGN_PARAGRAPH.LEFT or p.alignment is None:
                    alignment_str = 'LEFT'
                elif p.alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                    alignment_str = 'RIGHT'
                elif p.alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                    alignment_str = 'JUSTIFY'
                heading_alignments.append({'text': p.text[:60], 'alignment': alignment_str})
        return {'headings': heading_alignments, 'count': len(heading_alignments)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_sizes__5b4d6572a569badb37aa2fd25bab4b6e(env, config: dict):
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        sizes_pt = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    sizes_pt.append(run.font.size / 12700)
        style_size = None
        normal = doc.styles['Normal']
        if normal.font.size:
            style_size = normal.font.size / 12700
        return {'run_sizes': sizes_pt, 'style_size': style_size, 'total_runs': len(sizes_pt)}
    finally:
        os.unlink(tmp_path)

def get_docx_footer_page_numbers__9de26a1b58798be1fa2f4b6b3d9e9ec9(env, config: dict):
    """Check if the document has page numbers in footers."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        from docx.oxml.ns import qn
        doc = Document(tmp_path)
        has_page_number = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            footer_xml = footer._element.xml
            if 'PAGE' in footer_xml or 'w:fldChar' in footer_xml:
                has_page_number = True
                break
            for para in footer.paragraphs:
                if para.text.strip():
                    has_page_number = True
                    break
        return {'has_page_numbers': has_page_number}
    finally:
        os.unlink(tmp_path)

def get_docx_styled_merge__3faeea9d9342eb480cdf899557de9715(env, config: dict):
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found: ' + config.get('path', '')}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = []
        font_sizes = []
        font_names = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
                if run.font.name:
                    font_names.append(run.font.name)
        text = ' '.join(full_text)
        return {'text': text, 'font_sizes': list(set(font_sizes)) if font_sizes else [], 'font_names': list(set(font_names)) if font_names else []}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_first_para_alignment__f845858287fa4f02431e180c3025ad53(env, config: dict):
    """Get the alignment of the first paragraph (title) from the document."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) == 0:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        alignment = first_para.alignment
        is_centered = alignment == WD_ALIGN_PARAGRAPH.CENTER
        return {'alignment': str(alignment), 'is_centered': is_centered}
    finally:
        os.unlink(tmp_path)

def get_docx_reverse_merge__08f1db516ca13465963b37a9fdcaf66c(env, config: dict):
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found: ' + config.get('path', '')}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = []
        font_sizes = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
        text = ' '.join(full_text)
        first_150 = text[:150] if text else ''
        last_150 = text[-150:] if text else ''
        return {'text': text, 'first_150': first_150, 'last_150': last_150, 'font_sizes': list(set(font_sizes)) if font_sizes else []}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_selective_merge__1d49dccb31c44fca7886c9c825e7148f(env, config: dict):
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found: ' + config.get('path', '')}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = []
        font_sizes = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
        text = ' '.join(full_text)
        return {'text': text, 'font_sizes': list(set(font_sizes)) if font_sizes else [], 'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()])}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_para_alignment__4ed04700852dd5a18fbdc223aac7b1c3(env, config: dict):
    """Get alignment of a specific paragraph in a docx file."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        target_para_index = config.get('paragraph_index', 0)
        non_empty = [(i, p) for (i, p) in enumerate(doc.paragraphs) if p.text.strip()]
        if target_para_index >= len(non_empty):
            return {'error': f'Paragraph index {target_para_index} out of range'}
        (_, para) = non_empty[target_para_index]
        alignment = para.alignment
        alignment_map = {0: 'LEFT', 1: 'CENTER', 2: 'RIGHT', 3: 'JUSTIFY', None: 'LEFT'}
        alignment_int = alignment.value if hasattr(alignment, 'value') else alignment
        return {'alignment': alignment_map.get(alignment_int, str(alignment_int)), 'alignment_int': alignment_int}
    finally:
        os.unlink(tmp_path)

def get_docx_text__896d76d765c6621bc8feff2ba14e6be9(env, config: dict):
    """Read text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        text = '\n'.join((p.text for p in doc.paragraphs)).strip()
        return {'text': text}
    finally:
        os.unlink(tmp_path)

def get_power_settings_state__e637b95c625a88b3bb2cee1c77022b86(env, config: dict):
    """Get both idle-dim and idle-delay power settings."""
    try:
        idle_dim_result = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power idle-dim', timeout=30)
        idle_delay_result = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=30)
        idle_dim = idle_dim_result.get('output', '').strip() if isinstance(idle_dim_result, dict) else str(idle_dim_result).strip()
        idle_delay = idle_delay_result.get('output', '').strip() if isinstance(idle_delay_result, dict) else str(idle_delay_result).strip()
        return {'idle_dim': idle_dim, 'idle_delay': idle_delay}
    except Exception as e:
        return {'error': str(e)}

def get_docx_title_alignment__0a4a5cad7031961d3e6b497ec2c282f4(env, config: dict):
    """Get the alignment of the first paragraph (title) in the document."""
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    file_path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        alignment = first_para.paragraph_format.alignment
        alignment_map = {WD_PARAGRAPH_ALIGNMENT.CENTER: 'center', WD_PARAGRAPH_ALIGNMENT.LEFT: 'left', WD_PARAGRAPH_ALIGNMENT.RIGHT: 'right', WD_PARAGRAPH_ALIGNMENT.JUSTIFY: 'justify'}
        alignment_str = alignment_map.get(alignment, 'left')
        return {'alignment': alignment_str, 'title_text': first_para.text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_check_utils_py__597b85b971b86604637b8eea77967812(env, config: dict):
    """Check utils.py file existence, function definition, and correctness."""
    result = {}
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/project/utils.py')
        if file_bytes:
            content = file_bytes.decode('utf-8', errors='replace')
            result['file_exists'] = True
            result['has_add_function'] = 'def add' in content
        else:
            result['file_exists'] = False
            result['has_add_function'] = False
    except Exception:
        result['file_exists'] = False
        result['has_add_function'] = False
    try:
        check = env.controller.run_bash_script('cd /home/user/Desktop/project && python3 -c "from utils import add; print(add(2, 3))" 2>&1', timeout=10)
        output = check.get('output', '') if isinstance(check, dict) else str(check)
        result['func_output'] = output.strip()
    except Exception:
        result['func_output'] = 'ERROR'
    return result

def get_docx_font_info__2042c7daa4f47bea2c4fb2c5e4661307(env, config: dict):
    """Get font name and size from all runs in a docx file."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        fonts = set()
        sizes = set()
        for p in doc.paragraphs:
            if p.text.strip():
                for run in p.runs:
                    if run.text.strip():
                        if run.font.name:
                            fonts.add(run.font.name)
                        if run.font.size:
                            sizes.add(round(run.font.size / 12700, 1))
        return {'fonts': sorted(list(fonts)), 'sizes_pt': sorted(list(sizes))}
    finally:
        os.unlink(tmp_path)

def get_docx_text_formatting__92a351c44e937610849c0af1deb80b89(env, config: dict):
    """Get bold/italic/underline formatting for each non-empty paragraph."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        para_formatting = []
        for p in doc.paragraphs:
            if p.text.strip():
                runs_with_text = [run for run in p.runs if run.text.strip()]
                if runs_with_text:
                    is_bold = all((run.bold is True for run in runs_with_text))
                    is_italic = all((run.italic is True for run in runs_with_text))
                    is_underline = all((run.underline is not None and run.underline is not False for run in runs_with_text))
                else:
                    is_bold = False
                    is_italic = False
                    is_underline = False
                para_formatting.append({'bold': is_bold, 'italic': is_italic, 'underline': is_underline})
        return {'paragraphs': para_formatting}
    finally:
        os.unlink(tmp_path)

def get_header_italic__ee9130b48c4f13e968a7877f09e8d206(env, config: dict):
    """Check italic property of header cells."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cells_to_check = config.get('cells', ['B1', 'C1', 'D1', 'E1'])
        italic_states = {}
        for cell_ref in cells_to_check:
            cell = ws[cell_ref]
            italic_states[cell_ref] = bool(cell.font.italic)
        return {'italic_states': italic_states}
    finally:
        os.unlink(tmp_path)

def get_docx_italic_color__575ab439795c172e41189face7efb98a(env, config: dict):
    """Get font color of all italic runs in a docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_colors = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic:
                    color = run.font.color
                    if color and color.rgb:
                        italic_colors.append(str(color.rgb))
                    else:
                        italic_colors.append(None)
        return {'italic_colors': italic_colors, 'count': len(italic_colors)}
    finally:
        os.unlink(tmp_path)

def get_docx_first_para_font_size__0625ebacf8174a330d7b47b181788578(env, config: dict):
    """Get the font size of runs in the first paragraph."""
    from docx import Document
    from docx.shared import Pt
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        runs = first_para.runs
        if not runs:
            return {'font_size': None}
        sizes = []
        for run in runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
        if not sizes:
            return {'font_size': None}
        return {'font_size': sizes[0]}
    finally:
        os.unlink(tmp_path)

def get_background_color_state__0fc932ce1ad9456b6d290262a9551736(env, config: dict):
    """Download result and reference images from VM, return local paths for comparison."""
    result_bytes = env.controller.get_file(config['result_path'])
    ref_bytes = env.controller.get_file(config['reference_path'])
    if not result_bytes or not ref_bytes:
        return {'error': 'File not found'}
    tmp_result = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    tmp_result.write(result_bytes)
    tmp_result.close()
    tmp_ref = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    tmp_ref.write(ref_bytes)
    tmp_ref.close()
    return {'result_path': tmp_result.name, 'reference_path': tmp_ref.name}

def get_docx_para_font__ec9fd979c9f4ede06619498c1e9b1202(env, config: dict):
    """Get font name of runs in a specific paragraph of a docx file."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        target_para_index = config.get('paragraph_index', 0)
        non_empty = [(i, p) for (i, p) in enumerate(doc.paragraphs) if p.text.strip()]
        if not non_empty:
            return {'error': 'No non-empty paragraphs found'}
        if target_para_index >= len(non_empty) or target_para_index < -len(non_empty):
            return {'error': f'Paragraph index {target_para_index} out of range (total non-empty: {len(non_empty)})'}
        (_, para) = non_empty[target_para_index]
        total_chars = 0
        font_counts = {}
        for run in para.runs:
            run_len = len(run.text)
            total_chars += run_len
            fname = run.font.name or 'default'
            font_counts[fname] = font_counts.get(fname, 0) + run_len
        return {'total_chars': total_chars, 'font_counts': font_counts, 'dominant_font': max(font_counts, key=font_counts.get) if font_counts else 'unknown'}
    finally:
        os.unlink(tmp_path)

def get_row_hidden_state__18d397a77a413967bac86e92a1394add(env, config: dict):
    """Check hidden state of specific rows in the xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        rows_to_check = config.get('rows', [3])
        hidden_states = {}
        for row in rows_to_check:
            rd = ws.row_dimensions[row]
            hidden_states[str(row)] = bool(rd.hidden)
        return {'hidden_states': hidden_states}
    finally:
        os.unlink(tmp_path)

def get_docx_all_fonts__e754543a404f008975bc7bab115cbf09(env, config: dict):
    """Get all font names used in the document (paragraphs + table cells)."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        fonts = set()
        for para in doc.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts.add(run.font.name)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
        return {'fonts': list(fonts), 'font_count': len(fonts)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_tally_rows__89c74c0ac9693815f604c967708f8f82(env, config: dict):
    """Read all rows from the tally book xlsx and return as list of dicts."""
    import openpyxl
    path = config.get('path', '/home/user/Documents/Finance/tally_book.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        rows = []
        for row_idx in range(2, ws.max_row + 1):
            service = ws.cell(row=row_idx, column=1).value
            month = ws.cell(row=row_idx, column=2).value
            amount = ws.cell(row=row_idx, column=3).value
            rows.append({'row': row_idx, 'service': service, 'month': month, 'amount': amount})
        return {'rows': rows, 'max_row': ws.max_row}
    finally:
        os.unlink(tmp_path)

def get_docx_station_lines__dc2cb3918367317ad49987158bb0de7d(env, config: dict):
    """Get docx content and check for lines containing a specific station."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        station = config.get('station', 'SHS')
        station_lines = [l for l in lines if f',{station},' in l]
        non_station_lines = [l for l in lines if f',{station},' not in l]
        return {'total_lines': len(lines), 'station_lines_count': len(station_lines), 'non_station_lines_count': len(non_station_lines), 'has_station': len(station_lines) > 0, 'station': station}
    finally:
        os.unlink(tmp_path)

def get_check_first_line__35d6e56c2b031e6f068f1907dd34bebf(env, config: dict):
    """Get the first line of calculator.py to check if comment was added."""
    file_path = config.get('path', '/home/user/Desktop/calculator.py')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = content.splitlines()
    if not lines:
        return {'error': 'File is empty'}
    return {'first_line': lines[0].strip(), 'total_lines': len(lines)}

def get_docx_footer_page_numbers__806bfdf7c0ee49dcce08f91a2edad9e6(env, config: dict):
    """Check if the document has page numbers in footers."""
    from docx import Document
    from docx.oxml.ns import qn
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        has_page_number = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            footer_xml = footer._element.xml
            if 'PAGE' in footer_xml or 'page' in footer_xml.lower():
                has_page_number = True
                break
            for p in footer.paragraphs:
                for run in p.runs:
                    run_xml = run._element.xml
                    if 'PAGE' in run_xml:
                        has_page_number = True
                        break
                if has_page_number:
                    break
            if has_page_number:
                break
        return {'has_page_number': has_page_number}
    finally:
        os.unlink(tmp_path)

def get_docx_text__c53695cdc8b25b507dc6019962dcaf31(env, config: dict):
    """Read text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        text = '\n'.join((p.text for p in doc.paragraphs)).strip()
        return {'text': text}
    finally:
        os.unlink(tmp_path)

def get_writer_highlight_check__6b6e5048ebb3170359c47d1ce863bb95(env, config: dict):
    """Check which text spans have yellow highlighting in the document."""
    import tempfile
    import os
    file_path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from odf.opendocument import load
        from odf.style import Style, TextProperties
        from odf.text import P, Span
        doc = load(tmp_path)
        yellow_styles = set()
        styles = doc.automaticstyles.getElementsByType(Style)
        for style in styles:
            name = style.getAttribute('name')
            tps = style.getElementsByType(TextProperties)
            for tp in tps:
                bg = tp.getAttribute('backgroundcolor')
                if bg and bg.lower() in ('#ffff00', '#ffffff00', 'yellow'):
                    yellow_styles.add(name)
        body = doc.body
        paras = body.getElementsByType(P)
        highlighted_texts = []
        for (i, p) in enumerate(paras):
            spans = p.getElementsByType(Span)
            for s in spans:
                s_style = s.getAttribute('stylename')
                if s_style in yellow_styles:
                    text_content = ''
                    for node in s.childNodes:
                        text_content += str(node)
                    if text_content.strip():
                        highlighted_texts.append(text_content.strip())
        target_text = config.get('target_text', 'Study Team')
        all_text = ''
        for p in paras:
            p_text = ''
            for node in p.childNodes:
                p_text += str(node)
            all_text += p_text + '\n'
        target_count = all_text.count(target_text)
        highlighted_target_count = sum((1 for t in highlighted_texts if target_text in t or t in target_text))
        return {'highlighted_texts': highlighted_texts, 'highlighted_count': len(highlighted_texts), 'target_text': target_text, 'target_total_count': target_count, 'target_highlighted_count': highlighted_target_count}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_vtt_conversion_state__9ac6e58c407c5a10e99c8a824933ab2f(env, config: dict):
    """Check if subtitles were converted to WebVTT format."""
    target_path = config.get('path', '/home/user/captions.vtt')
    check_exists = env.controller.run_bash_script(f"test -f {target_path} && echo 'yes' || echo 'no'", timeout=10)
    file_exists = check_exists.get('output', '').strip() == 'yes'
    if not file_exists:
        return {'file_exists': False, 'has_webvtt_header': False, 'has_cues': False}
    check_content = env.controller.run_bash_script(f'head -20 {target_path} 2>/dev/null', timeout=10)
    content = check_content.get('output', '')
    has_webvtt_header = content.strip().startswith('WEBVTT')
    has_cues = '-->' in content
    return {'file_exists': file_exists, 'has_webvtt_header': has_webvtt_header, 'has_cues': has_cues}

def get_docx_page_numbers__18ec7a61b2f5504dc6e64dce31cd1e00(env, config: dict):
    """Check if document footers contain page numbers."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        has_page_numbers = False
        for section in doc.sections:
            footer = section.footer
            if footer:
                footer_xml = footer._element.xml
                if 'PAGE' in footer_xml:
                    has_page_numbers = True
                    break
                if 'w:fldSimple' in footer_xml:
                    has_page_numbers = True
                    break
                for p in footer.paragraphs:
                    p_xml = p._element.xml
                    if 'PAGE' in p_xml or 'w:fldChar' in p_xml:
                        has_page_numbers = True
                        break
                if has_page_numbers:
                    break
        return {'has_page_numbers': has_page_numbers}
    finally:
        os.unlink(tmp_path)

def get_user_password_home__6fbd0fb5ef1077b146e8ceece522a6b4(env, config: dict):
    """Get user existence, home directory, and password verification status."""
    username = config.get('username', 'charles')
    password = config.get('password', '')
    check_script = config.get('check_script', './check_password.sh')
    result_data = {'user_exists': False, 'home_correct': False, 'password_correct': False}
    try:
        user_info = env.controller.run_bash_script(f'getent passwd {username}', timeout=30)
        output = user_info.get('output', '').strip() if isinstance(user_info, dict) else str(user_info).strip()
        if not output:
            return result_data
        result_data['user_exists'] = True
        parts = output.split(':')
        if len(parts) >= 6:
            expected_home = config.get('expected_home', '/home/test1')
            result_data['home_correct'] = parts[5] == expected_home
        if password and check_script:
            pw_check = env.controller.run_bash_script(f'{check_script} "{username}" "{password}"', timeout=30)
            pw_output = pw_check.get('output', '').strip() if isinstance(pw_check, dict) else str(pw_check).strip()
            if 'success' in pw_output.lower() or pw_check.get('returncode', 1) == 0:
                result_data['password_correct'] = True
        return result_data
    except Exception as e:
        logger.error(f'Error checking user: {e}')
        return result_data

def get_writer_first_para_bold__974838092a7939215ca9c77d0b32d057(env, config: dict):
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        runs_bold = []
        for run in first_para.runs:
            runs_bold.append(run.bold is True)
        return {'runs_bold': runs_bold, 'total_runs': len(runs_bold), 'all_bold': all(runs_bold) if runs_bold else False}
    finally:
        os.unlink(tmp_path)

def get_docx_heading2_fonts__3ee05560ea0a8eae187c5c1b3f30c5d3(env, config: dict):
    """Get font names of all Heading 2 paragraphs in the docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        heading_fonts = []
        for p in doc.paragraphs:
            if p.style and p.style.name == 'Heading 2':
                fonts = set()
                for run in p.runs:
                    if run.font.name:
                        fonts.add(run.font.name)
                heading_fonts.append({'text': p.text[:60], 'fonts': list(fonts)})
        return {'headings': heading_fonts, 'count': len(heading_fonts)}
    finally:
        os.unlink(tmp_path)

def get_docx_title_text__ca26e505bd69174991058c87e74f2a98(env, config: dict):
    """Get the title (first non-empty paragraph) from a docx file on the VM."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        title = ''
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                title = text
                break
        return {'title': title}
    finally:
        os.unlink(tmp_path)

def get_docx_answer_line__12b6d7e1507f0e7212c8d9d4764a8a2b(env, config: dict):
    """Read Answer.docx and extract text after specified test header."""
    try:
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/Answer.docx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs]
            target_header = config.get('target_header', 'Grammar test 2:')
            result = {'paragraphs': paragraphs, 'answer_text': None}
            for (i, text) in enumerate(paragraphs):
                if text == target_header and i + 1 < len(paragraphs):
                    result['answer_text'] = paragraphs[i + 1]
                    break
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_italic_bold__7532d291294fb4f3bb8e25015109f986(env, config: dict):
    """Get bold status of all italic runs in a docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_bold_status = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic:
                    italic_bold_status.append(bool(run.bold))
        return {'italic_bold': italic_bold_status, 'count': len(italic_bold_status)}
    finally:
        os.unlink(tmp_path)

def get_docx_title_format__c874bb9d2ad63d6767d6a39e179946a5(env, config: dict):
    """Get the formatting of the first paragraph (title) in a docx file."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        alignment = first_para.alignment
        is_centered = alignment == WD_ALIGN_PARAGRAPH.CENTER
        all_bold = True
        has_runs = False
        for run in first_para.runs:
            if run.text.strip():
                has_runs = True
                if not run.bold:
                    all_bold = False
                    break
        if not has_runs:
            all_bold = False
        return {'is_centered': is_centered, 'is_bold': all_bold, 'title_text': first_para.text}
    finally:
        os.unlink(tmp_path)

def get_docx_first_para_font__53ba1671387b371fc528911be537db49(env, config: dict):
    """Get font name of the first paragraph's runs, with style inheritance fallback."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        target_para = None
        for p in doc.paragraphs:
            if p.text.strip():
                target_para = p
                break
        if target_para is None:
            return {'error': 'No non-empty paragraph found'}
        style_font = None
        if target_para.style and target_para.style.font and target_para.style.font.name:
            style_font = target_para.style.font.name
        fonts = []
        for run in target_para.runs:
            if run.font.name:
                fonts.append(run.font.name)
            elif style_font:
                fonts.append(style_font)
            else:
                fonts.append(None)
        return {'fonts': fonts, 'style_font': style_font, 'text_preview': target_para.text[:60]}
    finally:
        os.unlink(tmp_path)

def get_notification_dual_settings__0a0313fec36cad178ad94d7913e78201(env, config: dict):
    """Get both Do Not Disturb and Lock Screen Notifications settings."""
    try:
        dnd_result = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-banners', timeout=30)
        lock_result = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-in-lock-screen', timeout=30)
        dnd_value = dnd_result.get('output', '').strip() if isinstance(dnd_result, dict) else str(dnd_result).strip()
        lock_value = lock_result.get('output', '').strip() if isinstance(lock_result, dict) else str(lock_result).strip()
        return {'show_banners': dnd_value, 'show_in_lock_screen': lock_value}
    except Exception as e:
        return {'error': str(e)}

def get_writer_font_check__2a303ca90381d050b859a1a253c0648a(env, config: dict):
    """Get all font names used in the document text."""
    import tempfile
    import os
    file_path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from odf.opendocument import load
        from odf.style import Style, TextProperties
        from odf.text import P, Span
        doc = load(tmp_path)
        fonts_in_styles = {}
        styles = doc.automaticstyles.getElementsByType(Style)
        for style in styles:
            name = style.getAttribute('name')
            tps = style.getElementsByType(TextProperties)
            for tp in tps:
                fontname = tp.getAttribute('fontname')
                if fontname:
                    fonts_in_styles[name] = fontname
        body = doc.body
        paras = body.getElementsByType(P)
        used_fonts = set()
        for p in paras:
            p_style = p.getAttribute('stylename')
            if p_style and p_style in fonts_in_styles:
                used_fonts.add(fonts_in_styles[p_style])
            spans = p.getElementsByType(Span)
            for s in spans:
                s_style = s.getAttribute('stylename')
                if s_style and s_style in fonts_in_styles:
                    used_fonts.add(fonts_in_styles[s_style])
        default_styles = doc.styles.getElementsByType(Style)
        default_font = None
        for style in default_styles:
            name = style.getAttribute('name')
            if name in ('Standard', 'Default Paragraph Font'):
                tps = style.getElementsByType(TextProperties)
                for tp in tps:
                    fontname = tp.getAttribute('fontname')
                    if fontname:
                        default_font = fontname
        return {'used_fonts': sorted(list(used_fonts)), 'default_font': default_font, 'all_fonts_same': len(used_fonts) <= 1, 'single_font': list(used_fonts)[0] if len(used_fonts) == 1 else None}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_title_italic__088f05c0fd2c52a92bc3dc8f8cc6719c(env, config: dict):
    """Check if the first paragraph (title) runs are italic in a docx file."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        if not first_para.runs:
            return {'is_italic': False, 'text': first_para.text[:80]}
        all_italic = True
        any_run_checked = False
        for run in first_para.runs:
            if run.text.strip():
                any_run_checked = True
                if not run.font.italic:
                    all_italic = False
                    break
        if not any_run_checked:
            return {'is_italic': False, 'text': first_para.text[:80]}
        return {'is_italic': all_italic, 'text': first_para.text[:80]}
    finally:
        os.unlink(tmp_path)

def get_docx_footer_page_numbers__3205af8e0dea88590c54f355855d7cbf(env, config: dict):
    """Check if the document footers contain page number fields."""
    from docx import Document
    from lxml import etree
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        has_page_number = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            footer_xml = etree.tostring(footer._element, encoding='unicode')
            if 'PAGE' in footer_xml or 'fldChar' in footer_xml:
                has_page_number = True
                break
            for para in footer.paragraphs:
                if para.text.strip():
                    has_page_number = True
                    break
        return {'has_page_number': has_page_number}
    finally:
        os.unlink(tmp_path)

def get_contrast_comparison__a61095a992c03ac5e17b282d64e6f2d0(env, config: dict):
    """Fetch both original and edited images from VM, compute contrast and SSIM."""
    try:
        from PIL import Image
        import numpy as np
        edited_path = config.get('path', '')
        original_path = config.get('original_path', '')
        edited_bytes = env.controller.get_file(edited_path)
        if not edited_bytes:
            return {'error': f'Edited file not found: {edited_path}'}
        original_bytes = env.controller.get_file(original_path)
        if not original_bytes:
            return {'error': f'Original file not found: {original_path}'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_edited:
            tmp_edited.write(edited_bytes)
            tmp_edited_path = tmp_edited.name
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_original:
            tmp_original.write(original_bytes)
            tmp_original_path = tmp_original.name
        try:
            edited_img = Image.open(tmp_edited_path).convert('L')
            original_img = Image.open(tmp_original_path).convert('L')
            edited_arr = np.array(edited_img, dtype=np.float64)
            original_arr = np.array(original_img, dtype=np.float64)
            edited_contrast = float(np.std(edited_arr))
            original_contrast = float(np.std(original_arr))
            if edited_arr.shape != original_arr.shape:
                edited_img_resized = edited_img.resize(original_img.size, Image.LANCZOS)
                edited_arr_resized = np.array(edited_img_resized, dtype=np.float64)
            else:
                edited_arr_resized = edited_arr
            mse = float(np.mean((edited_arr_resized - original_arr) ** 2))
            similarity = max(0.0, 1.0 - mse / 65025.0)
            return {'original_contrast': original_contrast, 'edited_contrast': edited_contrast, 'similarity': similarity, 'contrast_increased': edited_contrast > original_contrast}
        finally:
            os.unlink(tmp_edited_path)
            os.unlink(tmp_original_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_italic_font_name__f2fb15859784992fd6cf15bd171d4b91(env, config: dict):
    """Get font names of all italic and non-italic runs in a docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_fonts = []
        non_italic_fonts = []
        for para in doc.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                font_name = run.font.name
                if run.italic:
                    italic_fonts.append(font_name)
                else:
                    non_italic_fonts.append(font_name)
        return {'italic_fonts': italic_fonts, 'italic_count': len(italic_fonts), 'non_italic_fonts': non_italic_fonts, 'non_italic_count': len(non_italic_fonts)}
    finally:
        os.unlink(tmp_path)

def get_settings_speed_blocksize__4b546ac4b7d86d470dd227a5c08930e6(env, config: dict):
    """Get GAME_SPEED and BLOCK_SIZE from settings.py."""
    path = config.get('path', '/home/user/Desktop/tetris/settings.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    speed_match = re.search('GAME_SPEED\\s*=\\s*(\\d+)', content)
    blocksize_match = re.search('BLOCK_SIZE\\s*=\\s*(\\d+)', content)
    return {'game_speed': int(speed_match.group(1)) if speed_match else None, 'block_size': int(blocksize_match.group(1)) if blocksize_match else None}

def get_docx_first_para_alignment__ae9bdf892a1bf6b8ab5ef93251d291e0(env, config: dict):
    """Get the alignment of the first paragraph in a DOCX file."""
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        alignment = first_para.paragraph_format.alignment
        alignment_map = {None: 'LEFT', WD_PARAGRAPH_ALIGNMENT.LEFT: 'LEFT', WD_PARAGRAPH_ALIGNMENT.CENTER: 'CENTER', WD_PARAGRAPH_ALIGNMENT.RIGHT: 'RIGHT', WD_PARAGRAPH_ALIGNMENT.JUSTIFY: 'JUSTIFY'}
        return {'alignment': alignment_map.get(alignment, str(alignment))}
    finally:
        os.unlink(tmp_path)

def get_docx_para_bold__7a1ffa51886994c5f3667d404d37d94d(env, config: dict):
    """Get bold formatting status of a specific paragraph in a docx file."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        target_para_index = config.get('paragraph_index', 0)
        non_empty = [(i, p) for (i, p) in enumerate(doc.paragraphs) if p.text.strip()]
        if target_para_index >= len(non_empty):
            return {'error': f'Paragraph index {target_para_index} out of range, only {len(non_empty)} non-empty paragraphs'}
        (_, para) = non_empty[target_para_index]
        total_chars = 0
        bold_chars = 0
        for run in para.runs:
            run_len = len(run.text)
            total_chars += run_len
            if run.font.bold:
                bold_chars += run_len
        return {'total_chars': total_chars, 'bold_chars': bold_chars, 'bold_ratio': bold_chars / total_chars if total_chars > 0 else 0.0}
    finally:
        os.unlink(tmp_path)

def get_docx_all_bold__aef5992d33ef93f956afe1a194ef0bce(env, config: dict):
    """Check if all text runs in the document are bold."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        total_runs = 0
        bold_runs = 0
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            for run in para.runs:
                if not run.text.strip():
                    continue
                total_runs += 1
                if run.bold is True:
                    bold_runs += 1
        return {'total_runs': total_runs, 'bold_runs': bold_runs, 'all_bold': total_runs > 0 and bold_runs == total_runs}
    finally:
        os.unlink(tmp_path)

def get_booking_form_state__cf1c94ef66b2087660e9bbd9fc7c9d7e(env, config: dict):
    """Get booking page URL and form input values."""
    result = {}
    try:
        tree = env.controller.get_accessibility_tree()
        url = ''
        for line in tree.split('\n'):
            if 'http' in line and ('outlook.office365.com' in line or 'mbta.com' in line):
                match = re.search('(https?://\\S+)', line)
                if match:
                    url = match.group(1)
                    break
        result['url'] = url
    except Exception:
        result['url'] = ''
    try:
        js_code = '\n        (function() {\n            var result = {};\n            try {\n                var nameInput = document.evaluate(\n                    "/html/body/div[2]/div/form/div[7]/div/div/div[1]/input[1]",\n                    document, null, XPathResult.FIRST_ORDERED_NODE_TYPE, null\n                ).singleNodeValue;\n                result[\'name\'] = nameInput ? nameInput.value : \'\';\n            } catch(e) { result[\'name\'] = \'\'; }\n            try {\n                var emailInput = document.evaluate(\n                    "/html/body/div[2]/div/form/div[7]/div/div/div[1]/input[2]",\n                    document, null, XPathResult.FIRST_ORDERED_NODE_TYPE, null\n                ).singleNodeValue;\n                result[\'mail\'] = emailInput ? emailInput.value : \'\';\n            } catch(e) { result[\'mail\'] = \'\'; }\n            return JSON.stringify(result);\n        })();\n        '
        js_result = env.controller.execute_js_command(js_code)
        if js_result and isinstance(js_result, str):
            import json
            form_data = json.loads(js_result)
            result['name'] = form_data.get('name', '')
            result['mail'] = form_data.get('mail', '')
        else:
            result['name'] = ''
            result['mail'] = ''
    except Exception:
        result['name'] = ''
        result['mail'] = ''
    return result

def get_docx_title_font__1b8a8292a0fd2eb13bf16e945356170a(env, config: dict):
    """Get the font name of the title paragraph in the docx document."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/H2O_Factsheet_WA.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) < 2:
            return {'error': 'Document has fewer than 2 paragraphs'}
        title_para = doc.paragraphs[1]
        title_text = title_para.text.strip()
        fonts = []
        for run in title_para.runs:
            font_name = run.font.name
            if font_name:
                fonts.append(font_name)
        return {'title_text': title_text, 'fonts': fonts, 'primary_font': fonts[0] if fonts else None}
    finally:
        os.unlink(tmp_path)

def get_docx_title_alignment__4f408a61496e9bff25fa5074f0f2e380(env, config: dict):
    """Get the alignment of the title paragraph in the docx document."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/H2O_Factsheet_WA.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) < 2:
            return {'error': 'Document has fewer than 2 paragraphs'}
        title_para = doc.paragraphs[1]
        title_text = title_para.text.strip()
        alignment = title_para.alignment
        alignment_map = {WD_ALIGN_PARAGRAPH.LEFT: 'left', WD_ALIGN_PARAGRAPH.CENTER: 'center', WD_ALIGN_PARAGRAPH.RIGHT: 'right', WD_ALIGN_PARAGRAPH.JUSTIFY: 'justify', None: 'left'}
        alignment_str = alignment_map.get(alignment, 'unknown')
        return {'title_text': title_text, 'alignment': alignment_str}
    finally:
        os.unlink(tmp_path)

def get_compose_recipients__9a5e3321e7daf20a2f35cb35407bc19e(env, config: dict):
    """Get compose window recipients from accessibility tree."""
    try:
        tree = env.controller.get_accessibility_tree()
        return {'tree_content': tree if isinstance(tree, str) else str(tree)}
    except Exception as e:
        return {'error': str(e)}

def get_writer_default_font__0f58b322b0ad4bcc09c937490ff9bfeb(env, config: dict):
    """Get the font of body text paragraphs in a docx file."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        fonts = {}
        for para in doc.paragraphs:
            if para.style.name and 'Heading' in para.style.name:
                continue
            if not para.text.strip():
                continue
            for run in para.runs:
                if run.text.strip():
                    font_name = run.font.name
                    if font_name:
                        fonts[font_name] = fonts.get(font_name, 0) + len(run.text)
        if not fonts:
            return {'font': None}
        most_common = max(fonts, key=fonts.get)
        return {'font': most_common, 'all_fonts': fonts}
    finally:
        os.unlink(tmp_path)

def get_docx_all_fonts__f0d9875ec4c08351c2fb33918a5793ea(env, config: dict):
    """Get all font names used in the document."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        fonts = set()
        for para in doc.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts.add(run.font.name)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
        return {'fonts': sorted(list(fonts))}
    finally:
        os.unlink(tmp_path)

def get_main_caption__909b46690f7fc746deb97149427ce266(env, config: dict):
    """Get the window caption from main.py."""
    path = config.get('path', '/home/user/Desktop/tetris/main.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    caption_match = re.search('set_caption\\(["\\\'](.+?)["\\\']\\)', content)
    return {'caption': caption_match.group(1) if caption_match else None}

def get_docx_first_para_alignment__58a0796c274d2d883482eaaa895cfb2f(env, config: dict):
    """Get the alignment of the first non-empty paragraph in a docx file."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        alignment_map = {WD_ALIGN_PARAGRAPH.LEFT: 'left', WD_ALIGN_PARAGRAPH.CENTER: 'center', WD_ALIGN_PARAGRAPH.RIGHT: 'right', WD_ALIGN_PARAGRAPH.JUSTIFY: 'justify', None: 'left'}
        for para in doc.paragraphs:
            if para.text.strip():
                alignment = alignment_map.get(para.alignment, 'unknown')
                return {'alignment': alignment, 'text': para.text.strip()[:80]}
        return {'error': 'No non-empty paragraphs found'}
    finally:
        os.unlink(tmp_path)

def get_writer_title_fontsize__22037d9b737ef0f8bf33646ac10692bd(env, config: dict):
    """Get the font size of the first non-empty paragraph (title) in the document."""
    import tempfile
    import os
    file_path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from odf.opendocument import load
        from odf.style import Style, TextProperties
        from odf.text import P, Span
        doc = load(tmp_path)
        style_fontsize = {}
        styles = doc.automaticstyles.getElementsByType(Style)
        for style in styles:
            name = style.getAttribute('name')
            tps = style.getElementsByType(TextProperties)
            for tp in tps:
                fontsize = tp.getAttribute('fontsize')
                if fontsize:
                    style_fontsize[name] = fontsize
        named_styles = doc.styles.getElementsByType(Style)
        for style in named_styles:
            name = style.getAttribute('name')
            tps = style.getElementsByType(TextProperties)
            for tp in tps:
                fontsize = tp.getAttribute('fontsize')
                if fontsize:
                    style_fontsize[name] = fontsize
        body = doc.body
        paras = body.getElementsByType(P)
        title_text = None
        title_fontsize = None
        title_para_style = None
        for p in paras:
            text_content = ''
            for node in p.childNodes:
                text_content += str(node)
            if text_content.strip():
                title_text = text_content.strip()
                title_para_style = p.getAttribute('stylename')
                if title_para_style and title_para_style in style_fontsize:
                    title_fontsize = style_fontsize[title_para_style]
                spans = p.getElementsByType(Span)
                for s in spans:
                    s_style = s.getAttribute('stylename')
                    if s_style and s_style in style_fontsize:
                        title_fontsize = style_fontsize[s_style]
                break
        return {'title_text': title_text, 'title_fontsize': title_fontsize, 'title_para_style': title_para_style}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_title_alignment__6fdb2c432a1af94ddff3a56f751c6296(env, config: dict):
    """Get the alignment of the title paragraph containing 'Graphemes' in the docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for p in doc.paragraphs:
            if 'Graphemes' in p.text and 'Sound Letter Patterns' in p.text:
                alignment = p.alignment
                if alignment is None:
                    alignment = p.style.paragraph_format.alignment
                is_centered = alignment == WD_ALIGN_PARAGRAPH.CENTER
                return {'alignment': str(alignment), 'is_centered': is_centered}
        return {'error': 'Title paragraph not found'}
    finally:
        os.unlink(tmp_path)

def get_budget_page_state__b76001106213215f01706a09eab507db(env, config: dict):
    """Get budget.com page URL and sort dropdown state from accessibility tree."""
    try:
        tree = env.controller.get_accessibility_tree()
        lines = tree.split('\n')
        url = ''
        for line in lines:
            m = re.search('https?://[^\\s\\\'">\\]\\)]+budget\\.com[^\\s\\\'">\\]\\)]*', line)
            if m:
                url = m.group(0)
                break
        sort_text = ''
        tree_lower = tree.lower()
        sort_pos = tree_lower.find('sort by')
        if sort_pos < 0:
            sort_pos = tree_lower.find('sort_by')
        if sort_pos < 0:
            sort_pos = tree_lower.find('sortby')
        if sort_pos >= 0:
            window = tree[sort_pos:sort_pos + 300]
            patterns = ['Price\\s*\\(Low to High\\)', 'Price\\s*\\(High to Low\\)', 'Number of Seats\\s*\\(High to Low\\)', 'Number of Seats\\s*\\(Low to High\\)', 'Recommended']
            for p in patterns:
                m = re.search(p, window, re.IGNORECASE)
                if m:
                    sort_text = m.group(0)
                    break
        return {'url': url, 'sort_text': sort_text}
    except Exception as e:
        return {'error': str(e)}

def get_check_uncomment__cd3236440d08e0a8882e8811427d01d3(env, config: dict):
    """Get the content of calculator.py to check if print statement is uncommented."""
    file_path = config.get('path', '/home/user/Desktop/calculator.py')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = content.splitlines()
    for (i, line) in enumerate(lines):
        stripped = line.strip()
        if 'Sorted array is:' in stripped:
            return {'line_content': stripped, 'line_number': i + 1, 'is_commented': stripped.startswith('#')}
    return {'error': 'Line with "Sorted array is:" not found'}

def get_docx_delimiter_check__0a313c15fc4fb10c48c753cea47e4527(env, config: dict):
    """Get docx content and check delimiter characters."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        old_delim = config.get('old_delimiter', ',')
        new_delim = config.get('new_delimiter', ';')
        old_count = sum((l.count(old_delim) for l in lines))
        new_count = sum((l.count(new_delim) for l in lines))
        first_line = lines[0] if lines else ''
        parts_old = first_line.split(old_delim)
        parts_new = first_line.split(new_delim)
        return {'total_lines': len(lines), 'old_delimiter_count': old_count, 'new_delimiter_count': new_count, 'first_line': first_line, 'parts_with_new_delim': len(parts_new)}
    finally:
        os.unlink(tmp_path)

def get_docx_font_check__c5ed5054f5dbce73ddb0483f71ff6f1a(env, config: dict):
    """Get font information from a docx document."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        default_font = None
        normal_style = doc.styles['Normal']
        if normal_style.font and normal_style.font.name:
            default_font = normal_style.font.name
        run_fonts = {}
        for p in doc.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    fn = r.font.name if r.font.name else 'inherited'
                    run_fonts[fn] = run_fonts.get(fn, 0) + 1
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        for r in p.runs:
                            if r.text.strip():
                                fn = r.font.name if r.font.name else 'inherited'
                                run_fonts[fn] = run_fonts.get(fn, 0) + 1
        return {'default_font': default_font, 'run_fonts': run_fonts}
    finally:
        os.unlink(tmp_path)

def get_docx_all_text_case__0c637a2f0d91975120dc8e0a76554911(env, config: dict):
    """Get all text from the docx and check if it's uppercase."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        all_text.append(cell.text)
        return {'paragraphs': all_text}
    finally:
        os.unlink(tmp_path)

def get_docx_heading_alignment__e69f89d4207a469ba2d4f747df580e7b(env, config: dict):
    """Get alignment of the 'Ontological Questions' heading paragraph."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        target_text = config.get('target_text', 'Ontological Questions')
        for p in doc.paragraphs:
            if target_text in p.text:
                alignment = p.paragraph_format.alignment
                if alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    align_str = 'center'
                elif alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                    align_str = 'right'
                elif alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                    align_str = 'justify'
                elif alignment == WD_ALIGN_PARAGRAPH.LEFT:
                    align_str = 'left'
                else:
                    align_str = 'left'
                return {'alignment': align_str, 'text': p.text.strip()}
        return {'error': f"Heading '{target_text}' not found"}
    finally:
        os.unlink(tmp_path)

def get_docx_page_orientation__7f5136b7f1dcbd0181f2d126bcac0caf(env, config: dict):
    """Get page orientation of a docx document."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        orientations = []
        for section in doc.sections:
            w = section.page_width
            h = section.page_height
            is_landscape = w > h
            orientations.append({'width': w, 'height': h, 'is_landscape': is_landscape})
        all_landscape = all((o['is_landscape'] for o in orientations)) if orientations else False
        any_landscape = any((o['is_landscape'] for o in orientations)) if orientations else False
        return {'all_landscape': all_landscape, 'any_landscape': any_landscape, 'section_count': len(orientations)}
    finally:
        os.unlink(tmp_path)

def get_settings_blue_color__be59b6debe1fbee8fc53696bdb2273b8(env, config: dict):
    """Get BLUE color constant from settings.py."""
    try:
        result = env.controller.run_bash_script('python3 -c "exec(open(\'/home/user/Desktop/snake/settings.py\').read()); print(BLUE)"', timeout=30)
        output = result.get('output', '').strip()
        if output.startswith('(') and output.endswith(')'):
            values = output.strip('()').split(',')
            if len(values) == 3:
                (r, g, b) = (int(values[0].strip()), int(values[1].strip()), int(values[2].strip()))
                return {'blue': [r, g, b], 'defined': True}
        return {'error': 'BLUE not found or invalid format', 'raw': output, 'defined': False}
    except Exception as e:
        return {'error': str(e), 'defined': False}

def get_docx_italic_status__bf5294aa1b721bbae785292729cb8020(env, config: dict):
    """Get italic status of all text runs in a docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        total_runs = 0
        italic_runs = 0
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    total_runs += 1
                    if run.italic:
                        italic_runs += 1
        return {'total_runs': total_runs, 'italic_runs': italic_runs, 'ratio': italic_runs / total_runs if total_runs > 0 else 0.0}
    finally:
        os.unlink(tmp_path)

def get_color_and_size_state__043ea3777ab786b45abdbd440171ffc3(env, config: dict):
    """Download result and reference images, return local paths and image dimensions."""
    from PIL import Image
    result_bytes = env.controller.get_file(config['result_path'])
    ref_bytes = env.controller.get_file(config['reference_path'])
    if not result_bytes or not ref_bytes:
        return {'error': 'File not found'}
    tmp_result = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    tmp_result.write(result_bytes)
    tmp_result.close()
    tmp_ref = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    tmp_ref.write(ref_bytes)
    tmp_ref.close()
    try:
        img = Image.open(tmp_result.name)
        (width, height) = img.size
    except Exception:
        (width, height) = (0, 0)
    return {'result_path': tmp_result.name, 'reference_path': tmp_ref.name, 'width': width, 'height': height}

def get_docx_first_para_italic__d7547f75d107bc34f4a878c280fffbc1(env, config: dict):
    """Check if all runs in the first paragraph are italic."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        runs = first_para.runs
        if not runs:
            return {'is_italic': False}
        all_italic = all((run.italic is True for run in runs))
        return {'is_italic': all_italic}
    finally:
        os.unlink(tmp_path)

def get_docx_heading_alignment__de2b396948b313eaa0067057e04e4e77(env, config: dict):
    """Get alignment of the first heading paragraph ('Example essay')."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for p in doc.paragraphs:
            if 'Example essay' in p.text:
                alignment = p.alignment
                if alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    align_str = 'CENTER'
                elif alignment == WD_ALIGN_PARAGRAPH.LEFT:
                    align_str = 'LEFT'
                elif alignment == WD_ALIGN_PARAGRAPH.RIGHT:
                    align_str = 'RIGHT'
                elif alignment == WD_ALIGN_PARAGRAPH.JUSTIFY:
                    align_str = 'JUSTIFY'
                elif alignment is None:
                    align_str = 'LEFT'
                else:
                    align_str = str(alignment)
                return {'text': p.text.strip(), 'alignment': align_str}
        return {'error': "Heading 'Example essay' not found"}
    finally:
        os.unlink(tmp_path)

def get_docx_title_formatting__6146cd6c34c0bc769a02e73ec8ef6f70(env, config: dict):
    """Get the formatting of the first paragraph (title) in the docx."""
    from docx import Document
    from docx.shared import Pt
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        is_bold = True
        font_sizes = []
        for run in first_para.runs:
            if run.text.strip():
                if not run.bold:
                    is_bold = False
                if run.font.size is not None:
                    font_sizes.append(run.font.size / 12700)
        return {'text': first_para.text, 'is_bold': is_bold, 'font_sizes': font_sizes}
    finally:
        os.unlink(tmp_path)

def get_docx_hint_line_spacing__7fbc318f6c5b72e9f18e5cf28491b9c0(env, config: dict):
    """Get line spacing of the five reading hint paragraphs (P15-P19)."""
    from docx import Document
    from docx.enum.text import WD_LINE_SPACING
    file_path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        hint_indices = []
        found_header = False
        for (i, p) in enumerate(doc.paragraphs):
            if 'Five hints which may be useful' in p.text:
                found_header = True
                continue
            if found_header and p.text.strip() and (len(hint_indices) < 5):
                hint_indices.append(i)
        spacings = []
        for idx in hint_indices:
            p = doc.paragraphs[idx]
            pf = p.paragraph_format
            spacing_rule = pf.line_spacing_rule
            spacing_val = pf.line_spacing
            if spacing_rule == WD_LINE_SPACING.ONE_POINT_FIVE_LINES:
                spacings.append(1.5)
            elif spacing_rule == WD_LINE_SPACING.DOUBLE:
                spacings.append(2.0)
            elif spacing_rule == WD_LINE_SPACING.SINGLE:
                spacings.append(1.0)
            elif spacing_val is not None:
                spacings.append(float(spacing_val))
            else:
                spacings.append(None)
        return {'spacings': spacings, 'hint_count': len(hint_indices)}
    finally:
        os.unlink(tmp_path)

def get_docx_footer_page_numbers__ad71c47cc02b102c764d5154ce0d73c0(env, config: dict):
    """Download docx from VM and check if footers contain page number fields."""
    import tempfile
    import os
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'has_page_numbers': False, 'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            has_page_numbers = False
            for section in doc.sections:
                footer = section.footer
                if footer and (not footer.is_linked_to_previous):
                    from lxml import etree
                    footer_xml = etree.tostring(footer._element, encoding='unicode')
                    if 'fldChar' in footer_xml or 'PAGE' in footer_xml or 'instrText' in footer_xml:
                        has_page_numbers = True
                        break
                even_footer = section.even_page_footer
                if even_footer:
                    even_xml = etree.tostring(even_footer._element, encoding='unicode')
                    if 'fldChar' in even_xml or 'PAGE' in even_xml or 'instrText' in even_xml:
                        has_page_numbers = True
                        break
            if not has_page_numbers:
                from docx.opc.constants import RELATIONSHIP_TYPE as RT
                for section in doc.sections:
                    footer = section.footer
                    if footer:
                        footer_xml = etree.tostring(footer._element, encoding='unicode')
                        if 'PAGE' in footer_xml:
                            has_page_numbers = True
                            break
            return {'has_page_numbers': has_page_numbers}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'has_page_numbers': False, 'error': str(e)}

def get_docx_all_fonts__2ae82910da68002280a1510e1ca61e99(env, config: dict):
    """Get all font names used in the docx document."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        fonts_found = []
        total_runs = 0
        matching_runs = 0
        target_font = config.get('target_font', 'Times New Roman')
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    total_runs += 1
                    font_name = run.font.name
                    if font_name:
                        fonts_found.append(font_name)
                        if font_name == target_font:
                            matching_runs += 1
                    elif para.style and para.style.font and para.style.font.name:
                        fonts_found.append(para.style.font.name)
                        if para.style.font.name == target_font:
                            matching_runs += 1
        default_font = None
        try:
            default_font = doc.styles['Normal'].font.name
        except Exception:
            pass
        return {'fonts': list(set(fonts_found)), 'default_font': default_font, 'total_runs': total_runs, 'matching_runs': matching_runs}
    finally:
        os.unlink(tmp_path)

def get_docx_text_replace_check__3df1073706339202301b753c7942ef7f(env, config: dict):
    """Check if '<add here>' was replaced with '(3)' in the document."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join([p.text for p in doc.paragraphs])
        has_add_here = '<add here>' in full_text
        has_replacement = '(3)' in full_text
        pennington_para_text = ''
        for p in doc.paragraphs:
            if 'Pennington' in p.text:
                pennington_para_text = p.text
                break
        return {'has_add_here': has_add_here, 'has_replacement': has_replacement, 'pennington_context': pennington_para_text[:200] if pennington_para_text else ''}
    finally:
        os.unlink(tmp_path)

def get_writer_title_alignment__3a67dac16aed3d74bffb9cb95c614713(env, config: dict):
    """Get the alignment of the first paragraph (title) in a docx file."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        first_para = doc.paragraphs[0]
        alignment = first_para.alignment
        align_map = {WD_ALIGN_PARAGRAPH.LEFT: 'left', WD_ALIGN_PARAGRAPH.CENTER: 'center', WD_ALIGN_PARAGRAPH.RIGHT: 'right', WD_ALIGN_PARAGRAPH.JUSTIFY: 'justify', None: 'left'}
        align_str = align_map.get(alignment, 'unknown')
        return {'alignment': align_str, 'text': first_para.text[:80]}
    finally:
        os.unlink(tmp_path)

def get_rows_hidden_state__cb276526eec2ec94231e7dedd20b5069(env, config: dict):
    """Check hidden state of specific rows in the xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        rows_to_check = config.get('rows', [32, 33])
        hidden_states = {}
        for row in rows_to_check:
            rd = ws.row_dimensions[row]
            hidden_states[str(row)] = bool(rd.hidden)
        return {'hidden_states': hidden_states}
    finally:
        os.unlink(tmp_path)

def get_docx_title_bold__94187b1698fde7f55883e20a5cb9822f(env, config: dict):
    """Get bold status of the title paragraph in the docx document."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/H2O_Factsheet_WA.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) < 2:
            return {'error': 'Document has fewer than 2 paragraphs'}
        title_para = doc.paragraphs[1]
        title_text = title_para.text.strip()
        all_bold = True
        run_count = len(title_para.runs)
        if run_count == 0:
            return {'error': 'Title paragraph has no runs', 'title_text': title_text}
        for run in title_para.runs:
            if not run.font.bold:
                all_bold = False
                break
        return {'title_text': title_text, 'all_bold': all_bold, 'run_count': run_count}
    finally:
        os.unlink(tmp_path)

def get_venv_check__3936cc1f2257a01f73d410787af3a332(env, config: dict):
    """Check if a Python virtual environment exists at the specified path."""
    venv_path = config.get('venv_path', '/home/user/myenv')
    results = {}
    cmd_result = env.controller.run_bash_script(f'test -d "{venv_path}" && echo "exists" || echo "missing"', timeout=10)
    output = cmd_result.get('output', '').strip() if isinstance(cmd_result, dict) else str(cmd_result).strip()
    results['venv_dir'] = output == 'exists'
    cmd_result = env.controller.run_bash_script(f'test -f "{venv_path}/bin/python" && echo "exists" || echo "missing"', timeout=10)
    output = cmd_result.get('output', '').strip() if isinstance(cmd_result, dict) else str(cmd_result).strip()
    results['python_binary'] = output == 'exists'
    cmd_result = env.controller.run_bash_script(f'test -f "{venv_path}/bin/pip" && echo "exists" || echo "missing"', timeout=10)
    output = cmd_result.get('output', '').strip() if isinstance(cmd_result, dict) else str(cmd_result).strip()
    results['pip_binary'] = output == 'exists'
    return results

def get_budget_reservation_info__a4a15f9cf85bf343423f5b15bc460e50(env, config: dict):
    """Get budget.com reservation page URL and pickup location from accessibility tree."""
    try:
        tree = env.controller.get_accessibility_tree()
        lines = tree.split('\n')
        url = ''
        for line in lines:
            m = re.search('https?://[^\\s\\\'">\\]\\)]+budget\\.com[^\\s\\\'">\\]\\)]*', line)
            if m:
                url = m.group(0)
                break
        tree_lower = tree.lower()
        has_boston_logan = 'boston logan' in tree_lower
        has_bos = bool(re.search('\\bBOS\\b', tree))
        return {'url': url, 'has_boston_logan': has_boston_logan, 'has_bos_code': has_bos}
    except Exception as e:
        return {'error': str(e)}

def get_settings_fps__6d1f5614df32aa00fba65a5a2ca43cb2(env, config: dict):
    """Get FPS value from settings.py."""
    try:
        result = env.controller.run_bash_script('python3 -c "exec(open(\'/home/user/Desktop/snake/settings.py\').read()); print(FPS)"', timeout=30)
        output = result.get('output', '').strip()
        fps = int(float(output))
        return {'fps': fps}
    except Exception as e:
        return {'error': str(e)}

def get_docx_footer_page_numbers__b8220411d07c53e298b6db0365205e12(env, config: dict):
    """Check if the docx has page numbers in any footer."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        has_page_number = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            for p in footer.paragraphs:
                for elem in p._element.iter():
                    tag = elem.tag
                    if tag.endswith('}fldSimple'):
                        instr = elem.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instr', '')
                        if 'PAGE' in instr.upper():
                            has_page_number = True
                    elif tag.endswith('}instrText'):
                        text = elem.text or ''
                        if 'PAGE' in text.upper():
                            has_page_number = True
        return {'has_page_number': has_page_number}
    finally:
        os.unlink(tmp_path)

def get_docx_default_font__fc8e11a362be23f66acd12dd6bff4956(env, config: dict):
    """Get the default font of the document."""
    import tempfile
    import os
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        default_font = None
        normal_style = doc.styles['Normal']
        if normal_style.font and normal_style.font.name:
            default_font = normal_style.font.name
        if not default_font:
            from docx.oxml.ns import qn
            doc_defaults = doc.styles.element.find(qn('w:docDefaults'))
            if doc_defaults is not None:
                rpr_default = doc_defaults.find(qn('w:rPrDefault'))
                if rpr_default is not None:
                    rpr = rpr_default.find(qn('w:rPr'))
                    if rpr is not None:
                        rfonts = rpr.find(qn('w:rFonts'))
                        if rfonts is not None:
                            default_font = rfonts.get(qn('w:ascii'))
        if not default_font:
            font_counts = {}
            for para in doc.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        font_counts[run.font.name] = font_counts.get(run.font.name, 0) + 1
            if font_counts:
                default_font = max(font_counts, key=font_counts.get)
        return {'default_font': default_font}
    finally:
        os.unlink(tmp_path)

def get_docx_heading_bold__0d3391a68f9d89ec2a40039cb5c06d2d(env, config: dict):
    """Check if the 'References' heading is bold."""
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for p in doc.paragraphs:
            if p.text.strip() == 'References':
                if not p.runs:
                    return {'text': p.text.strip(), 'is_bold': False, 'run_count': 0}
                all_bold = all((r.bold is True for r in p.runs))
                any_bold = any((r.bold is True for r in p.runs))
                return {'text': p.text.strip(), 'is_bold': all_bold, 'any_bold': any_bold, 'run_count': len(p.runs)}
        return {'error': "'References' heading not found"}
    finally:
        os.unlink(tmp_path)

def get_docx_answer_line__0636655694efccc23299c2ca2f236f1f(env, config: dict):
    """Read Answer.docx and extract text after specified test header."""
    try:
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/Answer.docx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            paragraphs = [p.text.strip() for p in doc.paragraphs]
            target_header = config.get('target_header', 'Grammar test 3:')
            result = {'paragraphs': paragraphs, 'answer_text': None}
            for (i, text) in enumerate(paragraphs):
                if text == target_header and i + 1 < len(paragraphs):
                    result['answer_text'] = paragraphs[i + 1]
                    break
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_text__72fc0161fa321ae4bf01e0be489f2375(env, config: dict):
    """Read text content from a docx file on the VM."""
    from docx import Document
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        text = '\n'.join((p.text for p in doc.paragraphs)).strip()
        return {'text': text}
    finally:
        os.unlink(tmp_path)

def get_docx_text_case__0a695088802a9cb1473c98bdb4ec1b4f(env, config: dict):
    """Get all text from a docx file for case checking."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_text = ''
        for para in doc.paragraphs:
            all_text += para.text + '\n'
        return {'text': all_text.strip()}
    finally:
        os.unlink(tmp_path)

def get_docx_first_line__2130fa798c3c78ff755136f88d8b3efd(env, config: dict):
    """Get the first paragraph of a docx document."""
    import tempfile
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        first_para_text = ''
        total_lines = 0
        for p in doc.paragraphs:
            if p.text.strip():
                if not first_para_text:
                    first_para_text = p.text.strip()
                total_lines += 1
        first_para = doc.paragraphs[0] if doc.paragraphs else None
        is_bold = False
        font_size = None
        if first_para and first_para.runs:
            is_bold = any((r.bold for r in first_para.runs))
            for r in first_para.runs:
                if r.font.size:
                    font_size = r.font.size
        return {'first_line': first_para_text, 'total_lines': total_lines, 'is_bold': is_bold, 'font_size': font_size}
    finally:
        os.unlink(tmp_path)

def get_docx_default_font__536a5616674c51027b6cacf28aa7ecd9(env, config: dict):
    """Get the default/Normal style font and all run fonts from a docx file."""
    from docx import Document
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        normal_style = doc.styles['Normal']
        default_font = normal_style.font.name
        run_fonts = set()
        for para in doc.paragraphs:
            for run in para.runs:
                if run.font.name:
                    run_fonts.add(run.font.name)
        return {'default_font': default_font, 'run_fonts': list(run_fonts)}
    finally:
        os.unlink(tmp_path)

def get_delta_miles_checkbox__f9b3359b12066a252ed8849a9876d835_qw35sft2_b6e19eb1(env, config: dict):
    """Check if 'Shop with Miles' checkbox is checked on delta.com.

    Returns {'checked': True/False}.
    """
    try:
        tree = env.controller.get_accessibility_tree()
        idx = tree.find('Shop with Miles')
        if idx == -1:
            lower = tree.lower()
            idx = lower.find('shop with miles')
            if idx == -1:
                return {'checked': False}
        context = tree[max(0, idx - 300):idx + 100]
        if re.search('\\bchecked\\b', context, re.IGNORECASE):
            return {'checked': True}
        return {'checked': False}
    except Exception as e:
        return {'error': str(e), 'checked': False}

def get_delta_form_state__45facd29054b2deec3804036a06322c9_qw35sft2_dd814e40(env, config: dict):
    """Get the Delta flight search form accessibility tree text."""
    try:
        tree = env.controller.get_accessibility_tree()
        tree_str = str(tree) if not isinstance(tree, str) else tree
        return {'tree_text': tree_str}
    except Exception as e:
        return {'error': str(e), 'tree_text': ''}

def get_delta_form_state__607589a24fcd7d8000928f466f0dd577_qw35sft2_503fe9b7(env, config: dict):
    """Get the Delta flight search form accessibility tree text."""
    try:
        tree = env.controller.get_accessibility_tree()
        tree_str = str(tree) if not isinstance(tree, str) else tree
        return {'tree_text': tree_str}
    except Exception as e:
        return {'error': str(e), 'tree_text': ''}

def get_delta_destination_field__696c0a5477f795f89ba4614ef9b530f9_qw35sft2_06ddfc62(env, config: dict):
    """Check if New York City Area (NYC) is set as the flight destination on delta.com.

    Returns {'destination_set': True/False, 'found_text': str}.
    """
    try:
        tree = env.controller.get_accessibility_tree()
        if re.search('\\bNYC\\b', tree) and 'New York' in tree:
            return {'destination_set': True, 'found_text': 'NYC New York'}
        if re.search('\\bNYC\\b', tree):
            return {'destination_set': True, 'found_text': 'NYC'}
        return {'destination_set': False, 'found_text': ''}
    except Exception as e:
        return {'error': str(e), 'destination_set': False}

def get_history_and_dnt__43a7a4d8fb3eedcedaf50f5a89a2ea93_qw35sft2_3461e6ae(env, config: dict):
    """Get YouTube history count and Do Not Track preference state."""
    import os
    import json
    import tempfile
    import sqlite3
    youtube_count = -1
    history_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/History')
    if history_bytes:
        with tempfile.NamedTemporaryFile(suffix='.sqlite', delete=False) as tmp:
            tmp.write(history_bytes)
            tmp_path = tmp.name
        try:
            conn = sqlite3.connect(tmp_path)
            cursor = conn.cursor()
            cursor.execute('SELECT url FROM urls')
            urls = [row[0] for row in cursor.fetchall()]
            conn.close()
            youtube_count = sum((1 for url in urls if 'youtube' in url.lower()))
        finally:
            os.unlink(tmp_path)
    do_not_track = False
    prefs_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
    if prefs_bytes:
        try:
            prefs = json.loads(prefs_bytes.decode('utf-8'))
            do_not_track = prefs.get('enable_do_not_track', False)
        except Exception:
            pass
    return {'youtube_count': youtube_count, 'do_not_track': do_not_track}

def get_font_and_dnt__baedb9f4aec1e7501564ae9f37c335ac_qw35sft2_30331902(env, config: dict):
    """Get Chrome default font size and Do Not Track setting from preferences."""
    result = env.controller.run_bash_script('find /root /home -name "Preferences" -path "*/google-chrome/Default/*" 2>/dev/null | head -1', timeout=10)
    if isinstance(result, dict):
        prefs_path = result.get('output', '').strip()
    else:
        prefs_path = str(result).strip()
    if not prefs_path:
        return {'error': 'Chrome Preferences file not found'}
    prefs_bytes = env.controller.get_file(prefs_path)
    if not prefs_bytes:
        return {'error': 'Cannot read Chrome Preferences file'}
    try:
        prefs = json.loads(prefs_bytes.decode('utf-8'))
        font_size = prefs.get('webkit', {}).get('webprefs', {}).get('default_font_size')
        do_not_track = prefs.get('enable_do_not_track', False)
        return {'font_size': font_size, 'do_not_track': do_not_track}
    except Exception as e:
        return {'error': str(e)}

def get_frame_at_3s__81411ca26f8f6e6edae115d08d9ab086_qw35sft2_cd74764a(env, config: dict):
    """Get user's frame.png and extract reference frame at 3s from fullvideo.mp4 for comparison."""
    user_frame_path = config.get('path', '/home/user/frame.png')
    video_path = config.get('video_path', '/home/user/fullvideo.mp4')
    ref_frame_path = '/tmp/ref_frame_81411ca26f8f6e6edae115d08d9ab086.png'
    check_result = env.controller.run_bash_script(f'test -f "{user_frame_path}" && echo EXISTS || echo MISSING', timeout=10)
    if not check_result or 'EXISTS' not in check_result.get('output', ''):
        return {'error': 'frame.png not found', 'frame_exists': False}
    ffmpeg_result = env.controller.run_bash_script(f'ffmpeg -y -i "{video_path}" -ss 3 -frames:v 1 "{ref_frame_path}" 2>&1 && echo FFMPEG_OK || echo FFMPEG_FAIL', timeout=30)
    if not ffmpeg_result or 'FFMPEG_OK' not in ffmpeg_result.get('output', ''):
        return {'error': 'Failed to extract reference frame from video', 'frame_exists': True}
    user_bytes = env.controller.get_file(user_frame_path)
    ref_bytes = env.controller.get_file(ref_frame_path)
    if not user_bytes:
        return {'error': 'Could not read user frame.png', 'frame_exists': True}
    if not ref_bytes:
        return {'error': 'Could not read extracted reference frame', 'frame_exists': True}
    return {'frame_exists': True, 'user_frame': base64.b64encode(user_bytes).decode('utf-8'), 'ref_frame': base64.b64encode(ref_bytes).decode('utf-8')}

def get_freeze_panes_state__10368a4827622b87edd2c56b2f249e0f_qw35sft2_2e501aba(env, config: dict):
    """Get the freeze_panes setting from the LibreOffice Calc xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Freeze_row_column.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        freeze_panes = ws.freeze_panes
        return {'freeze_panes': str(freeze_panes) if freeze_panes else None}
    finally:
        os.unlink(tmp_path)

def get_salesrep_jan_total__5238964d8e657c42c8a059231010b871_qw35sft2_67ceaa27(env, config: dict):
    """Read cell B12 from SalesRep.xlsx to check the January total."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws['B12'].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_income_gross_with_total__9bcff5517fb7493e61233c0a423569f9_qw35sft2_58941de0(env, config: dict):
    """Read Gross Profit values J2:J10 and total row J11 from IncomeStatement2.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/IncomeStatement2.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        gross_profit = []
        for row in range(2, 11):
            val = ws.cell(row=row, column=10).value
            gross_profit.append(int(val) if val is not None else None)
        total_val = ws.cell(row=11, column=10).value
        total = int(total_val) if total_val is not None else None
        return {'gross_profit': gross_profit, 'total': total}
    finally:
        os.unlink(tmp_path)

def get_employee_ages_avg__151640eb3d43a7eb22c551550e103953_qw35sft2_74381784(env, config: dict):
    """Get DOB from column C, ages from column D, and the AVERAGE summary value from E2."""
    file_path = config.get('path', '/home/user/Employee_Age_By_Birthday.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        dob_values = []
        age_values = []
        for row in range(2, 30):
            dob = ws.cell(row=row, column=3).value
            age = ws.cell(row=row, column=4).value
            dob_values.append(dob)
            age_values.append(age)
        e2_value = ws.cell(row=2, column=5).value
        return {'dob_values': [str(d.date()) if hasattr(d, 'date') else d for d in dob_values], 'age_values': age_values, 'e2_value': float(e2_value) if e2_value is not None else None}
    finally:
        os.unlink(tmp_path)

def get_vlookup_f2_single__dba068db6ac4d383c892aae7e4d7fc62_qw35sft2_44448389(env, config: dict):
    """Get the value of cell F2 from VLOOKUP_Fill_the_form.xlsx."""
    file_bytes = env.controller.get_file('/home/user/VLOOKUP_Fill_the_form.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'F2': ws['F2'].value}
    finally:
        os.unlink(tmp_path)

def get_total_label_and_jan__4e23e64ce56ffa608b1e8fd2866f5ae0_qw35sft2_58708e5a(env, config: dict):
    """Read only A12 (label) and B12 (January total) from SalesRep.xlsx."""
    path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'label': ws['A12'].value, 'jan': ws['B12'].value}
    finally:
        os.unlink(tmp_path)

def get_header_bold_format__988aa3731e056dbe2f61a28b2135e9ae_qw35sft2_7afefb80(env, config: dict):
    """Get bold formatting of header cells A1 and B1."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Set_Decimal_Separator_Dot.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        a1_bold = ws['A1'].font.bold if ws['A1'].font else False
        b1_bold = ws['B1'].font.bold if ws['B1'].font else False
        return {'A1_bold': bool(a1_bold), 'B1_bold': bool(b1_bold)}
    finally:
        os.unlink(tmp_path)

def get_monthly_totals_row__aff1e90d6f0bac6579c8d1279e6c6c3e_qw35sft2_721b5769(env, config: dict):
    """Read cells B24, C24, D24, E24 from the OrderId_Month_Chart.xlsx file to check the totals row."""
    file_path = config.get('path', '/home/user/OrderId_Month_Chart.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'label': ws['B24'].value, 'jan_total': ws['C24'].value, 'feb_total': ws['D24'].value, 'mar_total': ws['E24'].value}
    finally:
        os.unlink(tmp_path)

def get_weekly_sales_profit_total_row__50f7148385ea6d504119cd2e40f7d1be_qw35sft2_e95480fb(env, config: dict):
    """Read Profit header, A12 label, and D12 total from WeeklySales.xlsx."""
    file_bytes = env.controller.get_file('/home/user/WeeklySales.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        profit_values = [ws.cell(row=r, column=4).value for r in range(2, 12)]
        return {'d1_header': ws['D1'].value, 'a12_label': ws['A12'].value, 'd12_total': ws['D12'].value, 'profit_values': profit_values}
    finally:
        os.unlink(tmp_path)

def get_employee_split_sorted__5d1c2e4438b4aac3fb5db710d29300f5_qw35sft2_bcc33200(env, config: dict):
    """Read Employee_Roles_and_Ranks.xlsx and return B/C/D column data for sort verification."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Employee_Roles_and_Ranks.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        rows = []
        for row in range(2, ws.max_row + 1):
            b = ws.cell(row=row, column=2).value
            c = ws.cell(row=row, column=3).value
            d = ws.cell(row=row, column=4).value
            rows.append({'first': str(b).strip() if b is not None else None, 'last': str(c).strip() if c is not None else None, 'rank': str(d).strip() if d is not None else None})
        return {'rows': rows, 'row_count': len(rows)}
    finally:
        os.unlink(tmp_path)

def get_period_rate_max_in_d1__14298c6976685b9e8d10b60e8490521a_qw35sft2_14e7047e(env, config: dict):
    """Read Period Rate header, C20 font color, and D1 value from PeriodRate.xlsx."""
    path = config.get('path', '/home/user/PeriodRate.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        c20_color = ''
        try:
            cell = ws['C20']
            if cell.font and cell.font.color and (cell.font.color.type == 'rgb'):
                c20_color = cell.font.color.rgb[-6:].lower()
        except Exception:
            pass
        return {'C1': ws['C1'].value, 'C20_font_color': c20_color, 'D1_value': ws['D1'].value}
    finally:
        os.unlink(tmp_path)

def get_maturity_sorted__6a0b6ecdb454d8389c50fe9e251b1580_qw35sft2_e71a2e21(env, config: dict):
    """Read MaturityDate.xlsx and return header C1 plus sorted maturity dates C2:C10."""
    path = config.get('path', '/home/user/MaturityDate.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header_c = ws['C1'].value
        if isinstance(header_c, str):
            header_c = header_c.strip()
        maturity_dates = []
        for row in range(2, 11):
            val = ws.cell(row=row, column=3).value
            if val is None:
                date_str = None
            elif hasattr(val, 'date'):
                date_str = val.date().isoformat()
            elif isinstance(val, (int, float)):
                base = date(1899, 12, 30)
                date_str = (base + timedelta(days=int(val))).isoformat()
            else:
                date_str = str(val).strip()
            maturity_dates.append(date_str)
        non_null = [d for d in maturity_dates if d is not None]
        is_sorted = non_null == sorted(non_null)
        return {'header_c': header_c, 'maturity_dates': maturity_dates, 'is_sorted': is_sorted, 'c2_value': maturity_dates[0] if maturity_dates else None, 'c10_value': maturity_dates[-1] if maturity_dates else None}
    finally:
        os.unlink(tmp_path)

def get_salesrep_label__d6d9e5e1d2ca3d967d1969fd1a1e0c2c_qw35sft2_9789eb01(env, config: dict):
    """Read cell A12 from SalesRep.xlsx to check for a 'Total' row label."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws['A12'].value
        return {'label': value}
    finally:
        os.unlink(tmp_path)

def get_seqno_and_total_row__f9ba38ec0e00ec28c9c0ef5057c79916_qw35sft2_c93bd6f3(env, config: dict):
    """Read Seq No. column (B2:B29), label in D30, and total value in E30."""
    path = config.get('path', '/home/user/Order_Sales_Serial#.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        seq_nos = [ws.cell(row=r, column=2).value for r in range(2, 30)]
        d30 = ws.cell(row=30, column=4).value
        e30 = ws.cell(row=30, column=5).value
        return {'seq_nos': seq_nos, 'd30': d30, 'e30': e30}
    finally:
        os.unlink(tmp_path)

def get_employee_ages__add589b113a1be1eec737c1ab22661fd_qw35sft2_b15815f6(env, config: dict):
    """Get DOB values from column C and age values from column D of Employee_Age_By_Birthday.xlsx."""
    file_path = config.get('path', '/home/user/Employee_Age_By_Birthday.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        dob_values = []
        age_values = []
        for row in range(2, 30):
            dob = ws.cell(row=row, column=3).value
            age = ws.cell(row=row, column=4).value
            dob_values.append(dob)
            age_values.append(age)
        return {'dob_values': [str(d.date()) if hasattr(d, 'date') else d for d in dob_values], 'age_values': age_values}
    finally:
        os.unlink(tmp_path)

def get_vlookup_f2_f4_rows__22ff441442432fb67bd17ef3eb97a292_qw35sft2_21d8f812(env, config: dict):
    """Get values of cells F2, F3, F4 from VLOOKUP_Fill_the_form.xlsx."""
    file_bytes = env.controller.get_file('/home/user/VLOOKUP_Fill_the_form.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'F2': ws['F2'].value, 'F3': ws['F3'].value, 'F4': ws['F4'].value}
    finally:
        os.unlink(tmp_path)

def get_income_net_sales_gross__c6e3aa9fd1cbfd34b4473b9e8f31e349_qw35sft2_7e167ef8(env, config: dict):
    """Read Net Sales (E2:E10) and Gross Profit (J2:J10) from IncomeStatement2.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/IncomeStatement2.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        net_sales = []
        gross_profit = []
        for row in range(2, 11):
            e_val = ws.cell(row=row, column=5).value
            j_val = ws.cell(row=row, column=10).value
            net_sales.append(int(e_val) if e_val is not None else None)
            gross_profit.append(int(j_val) if j_val is not None else None)
        return {'net_sales': net_sales, 'gross_profit': gross_profit}
    finally:
        os.unlink(tmp_path)

def get_total_row_state__5ac69261c8fb294b16f18049ece06ce9_qw35sft2_313652f6(env, config: dict):
    """Read the Total row (row 12) from SalesRep.xlsx: label in A12 and sums B12:G12."""
    path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        label = ws['A12'].value
        jan = ws['B12'].value
        feb = ws['C12'].value
        mar = ws['D12'].value
        apr = ws['E12'].value
        may = ws['F12'].value
        jun = ws['G12'].value
        return {'label': label, 'jan': jan, 'feb': feb, 'mar': mar, 'apr': apr, 'may': may, 'jun': jun}
    finally:
        os.unlink(tmp_path)

def get_row_hidden_state__93bd58adf830a704a8880938e35d28c4_qw35sft2_83a00ea9(env, config: dict):
    """Read xlsx and return hidden state for row 3 (5/7/2022, the first #N/A row)."""
    path = config.get('path', '/home/user/Date_Budget_Variance_HideNA.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        target_row = config.get('target_row', 3)
        hidden = False
        if target_row in ws.row_dimensions:
            hidden = bool(ws.row_dimensions[target_row].hidden)
        return {'hidden': hidden, 'target_row': target_row}
    finally:
        os.unlink(tmp_path)

def get_weekly_sales_sorted_by_profit__99a00a677b41fb8db78ec79d5381f0b2_qw35sft2_eae34b45(env, config: dict):
    """Read week names, profit header, and profit values from WeeklySales.xlsx to verify sort order."""
    file_bytes = env.controller.get_file('/home/user/WeeklySales.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        weeks = [ws.cell(row=r, column=1).value for r in range(2, 12)]
        profit_values = [ws.cell(row=r, column=4).value for r in range(2, 12)]
        return {'d1_header': ws['D1'].value, 'weeks': weeks, 'profit_values': profit_values}
    finally:
        os.unlink(tmp_path)

def get_employee_split_bold_header__fc565b338fbb3db18703aef2eaa623cb_qw35sft2_4d4d17f1(env, config: dict):
    """Read Employee_Roles_and_Ranks.xlsx and return split data + header bold formatting."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Employee_Roles_and_Ranks.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb_vals = openpyxl.load_workbook(tmp_path, data_only=True)
        wb_fmt = openpyxl.load_workbook(tmp_path, data_only=False)
        ws_vals = wb_vals.worksheets[0]
        ws_fmt = wb_fmt.worksheets[0]
        b2 = ws_vals.cell(row=2, column=2).value
        c2 = ws_vals.cell(row=2, column=3).value
        d2 = ws_vals.cell(row=2, column=4).value
        b_filled = sum((1 for row in range(2, ws_vals.max_row + 1) if ws_vals.cell(row=row, column=2).value is not None))
        b1_bold = ws_fmt.cell(row=1, column=2).font.bold if ws_fmt.cell(row=1, column=2).font else False
        c1_bold = ws_fmt.cell(row=1, column=3).font.bold if ws_fmt.cell(row=1, column=3).font else False
        d1_bold = ws_fmt.cell(row=1, column=4).font.bold if ws_fmt.cell(row=1, column=4).font else False
        return {'b2': str(b2).strip() if b2 is not None else None, 'c2': str(c2).strip() if c2 is not None else None, 'd2': str(d2).strip() if d2 is not None else None, 'b_filled_count': b_filled, 'b1_bold': bool(b1_bold), 'c1_bold': bool(c1_bold), 'd1_bold': bool(d1_bold)}
    finally:
        os.unlink(tmp_path)

def get_ramp_accel_diff__e64652ba3bf9a74b525e231bf3f391d2_qw35sft2_6edf4814(env, config: dict):
    """Read Acceleration Difference column (E = B - D) and B30 from RampUpAndDown.xlsx."""
    import tempfile
    import os
    import openpyxl
    file_bytes = env.controller.get_file(config.get('path', '/home/user/RampUpAndDown.xlsx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        e1_val = ws.cell(row=1, column=5).value
        e2_val = ws.cell(row=2, column=5).value
        e30_val = ws.cell(row=30, column=5).value
        b30_val = ws['B30'].value
        d30_val = ws['D30'].value
        return {'e1': str(e1_val).strip() if e1_val is not None else None, 'e2': e2_val, 'e30': e30_val, 'b30': b30_val, 'd30': d30_val}
    finally:
        os.unlink(tmp_path)

def get_period_rate_sum_c26__fd409a83b18f09cf45613c5173377d4d_qw35sft2_ec0397ab(env, config: dict):
    """Read Period Rate header, max-row font color (C20), and sum cell (C26) from PeriodRate.xlsx."""
    path = config.get('path', '/home/user/PeriodRate.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell_c20 = ws['C20']
        font_color = ''
        try:
            if cell_c20.font and cell_c20.font.color and (cell_c20.font.color.type == 'rgb'):
                font_color = cell_c20.font.color.rgb[-6:].lower()
        except Exception:
            pass
        return {'C1': ws['C1'].value, 'C20_value': cell_c20.value, 'C20_font_color': font_color, 'C26_value': ws['C26'].value}
    finally:
        os.unlink(tmp_path)

def get_maturity_total__461252d0c8044977d30ecc240f2dc9cd_qw35sft2_08c41667(env, config: dict):
    """Read MaturityDate.xlsx and return C1 header, A11 label, and B11 total loan days."""
    path = config.get('path', '/home/user/MaturityDate.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header_c = ws['C1'].value
        if isinstance(header_c, str):
            header_c = header_c.strip()
        a11_val = ws['A11'].value
        if isinstance(a11_val, str):
            a11_val = a11_val.strip()
        b11_val = ws['B11'].value
        if isinstance(b11_val, float) and b11_val == int(b11_val):
            b11_val = int(b11_val)
        return {'header_c': header_c, 'a11': a11_val, 'b11': b11_val}
    finally:
        os.unlink(tmp_path)

def get_freeze_panes_state__ae87df329e6b7ba3f255d9a2acdfeda6_qw35sft2_8fbda1a7(env, config: dict):
    """Get the freeze_panes setting from the LibreOffice Calc xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Freeze_row_column.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        freeze_panes = ws.freeze_panes
        return {'freeze_panes': str(freeze_panes) if freeze_panes else None}
    finally:
        os.unlink(tmp_path)

def get_seqno_and_sales_sum__00713006369c01f94f764c7eb6008543_qw35sft2_bdbb170f(env, config: dict):
    """Read Seq No. column (B2:B29) and total sales cell E30 from the xlsx file."""
    path = config.get('path', '/home/user/Order_Sales_Serial#.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        seq_nos = [ws.cell(row=r, column=2).value for r in range(2, 30)]
        e30 = ws.cell(row=30, column=5).value
        return {'seq_nos': seq_nos, 'e30': e30}
    finally:
        os.unlink(tmp_path)

def get_salesrep_last3_totals__ce77d41e3a7b3efa48e4655c22c8aa27_qw35sft2_dac8c6be(env, config: dict):
    """Read cells E12:G12 from SalesRep.xlsx to check Apr, May, Jun monthly totals."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'apr': ws['E12'].value, 'may': ws['F12'].value, 'jun': ws['G12'].value}
    finally:
        os.unlink(tmp_path)

def get_vlookup_f2_f12_all__202132c1158925d79d1ba222174d8f66_qw35sft2_e41cb20f(env, config: dict):
    """Get values of all officer name cells F2:F12 from VLOOKUP_Fill_the_form.xlsx."""
    file_bytes = env.controller.get_file('/home/user/VLOOKUP_Fill_the_form.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'F2': ws['F2'].value, 'F3': ws['F3'].value, 'F4': ws['F4'].value, 'F5': ws['F5'].value, 'F6': ws['F6'].value, 'F7': ws['F7'].value, 'F8': ws['F8'].value, 'F9': ws['F9'].value, 'F10': ws['F10'].value, 'F11': ws['F11'].value, 'F12': ws['F12'].value}
    finally:
        os.unlink(tmp_path)

def get_income_net_sales_and_cost__43eee44b990a73c4b5a6e2bc138833a7_qw35sft2_975e4d4c(env, config: dict):
    """Read Net Sales (E2:E10) and Total Cost (I2:I10) from IncomeStatement2.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/IncomeStatement2.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        net_sales = []
        total_cost = []
        for row in range(2, 11):
            e_val = ws.cell(row=row, column=5).value
            i_val = ws.cell(row=row, column=9).value
            net_sales.append(int(e_val) if e_val is not None else None)
            total_cost.append(int(i_val) if i_val is not None else None)
        return {'net_sales': net_sales, 'total_cost': total_cost}
    finally:
        os.unlink(tmp_path)

def get_strikethrough_and_transition__b374fc122f888120a74e40bee1da5133_qw35sft2_e3f3d361(env, config: dict):
    """Get strikethrough state for first two Finance Meetings items and slide 5 transition type."""
    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        file_bytes = env.controller.get_file('/home/user/Desktop/New_Club_Spring_2018_Training.pptx')
        if not file_bytes:
            return {'error': 'file not found'}
        prs = Presentation(io.BytesIO(file_bytes))
        strike_state = {}
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue
                    runs = [r for r in para.runs if r.text.strip()]
                    if runs:
                        has_strike = all((r.font.strike is True for r in runs))
                        strike_state[para_text] = has_strike
        slide5 = prs.slides[4]
        transition_type = None
        try:
            spPr = slide5._element
            mc_elem = spPr.find(qn('p:transition'))
            if mc_elem is not None:
                for child in mc_elem:
                    local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    transition_type = local
                    break
                if transition_type is None:
                    transition_type = 'unknown'
            else:
                transition_type = 'none'
        except Exception:
            transition_type = 'none'
        return {'strike_state': strike_state, 'transition_type': transition_type}
    except Exception as e:
        return {'error': str(e)}

def get_title_font_props__d1408966fdb4fd77bea2fd3e21b1ee11_qw35sft2_60a6ab82(env, config: dict):
    """Get font properties of the first text shape on slide 1 of the pptx file."""
    import tempfile, os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/39_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                tf = shape.text_frame
                bold_vals, size_vals, underline_vals, italic_vals = ([], [], [], [])
                for para in tf.paragraphs:
                    for run in para.runs:
                        bold_vals.append(run.font.bold)
                        size_vals.append(run.font.size)
                        underline_vals.append(run.font.underline)
                        italic_vals.append(run.font.italic)
                if bold_vals:
                    size_emu = size_vals[0]
                    size_pt = round(size_emu / 12700, 1) if size_emu else None
                    return {'bold': bold_vals[0], 'size_pt': size_pt, 'underline': underline_vals[0], 'italic': italic_vals[0], 'text': tf.text.strip()[:60]}
        return {'error': 'No text shapes found on slide 1'}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_strikethrough_and_italic__be4183a9c6982ba683532f14bc50d4bc_qw35sft2_3f4cf417(env, config: dict):
    """Get strikethrough state for two Finance Meetings items and italic state for third."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file('/home/user/Desktop/New_Club_Spring_2018_Training.pptx')
        if not file_bytes:
            return {'error': 'file not found'}
        prs = Presentation(io.BytesIO(file_bytes))
        slide = prs.slides[4]
        state = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue
                runs = [r for r in para.runs if r.text.strip()]
                if runs:
                    has_strike = all((r.font.strike is True for r in runs))
                    is_italic = all((r.font.italic is True for r in runs))
                    state[para_text] = {'strike': has_strike, 'italic': is_italic}
        return state
    except Exception as e:
        return {'error': str(e)}

def get_picture_heights__eff65ebb8ed6d102db81f31da1a823ed_qw35sft2_c0347b22(env, config: dict):
    """Get picture heights (cm) on slides 3, 4, 6 from 30_1.pptx."""
    import tempfile, os
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    file_bytes = env.controller.get_file('/home/user/Desktop/30_1.pptx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide3 = prs.slides[2]
        slide3_target = None
        for shape in slide3.shapes:
            if shape.shape_type == 5:
                if slide3_target is None or shape.height > slide3_target.height:
                    slide3_target = shape
        slide4 = prs.slides[3]
        slide4_target = None
        for shape in slide4.shapes:
            if shape.shape_type == 13 and shape.height > 5 * 360000:
                if slide4_target is None or shape.height > slide4_target.height:
                    slide4_target = shape
        slide6 = prs.slides[5]
        slide6_target = None
        for shape in slide6.shapes:
            if shape.shape_type == 13:
                if slide6_target is None or shape.width > slide6_target.width:
                    slide6_target = shape
        return {'slide3_height_cm': round(slide3_target.height / 360000.0, 4) if slide3_target else None, 'slide4_height_cm': round(slide4_target.height / 360000.0, 4) if slide4_target else None, 'slide6_height_cm': round(slide6_target.height / 360000.0, 4) if slide6_target else None}
    finally:
        os.unlink(tmp_path)

def get_title_font_props__f083039072ffd98eec2d101bb28501a7_qw35sft2_a86e06fe(env, config: dict):
    """Get font properties of the first text shape on slide 1 of the pptx file."""
    import tempfile, os
    from pptx import Presentation
    file_path = config.get('path', '/home/user/Desktop/39_2.pptx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        slide = prs.slides[0]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                tf = shape.text_frame
                bold_vals, size_vals, underline_vals, italic_vals = ([], [], [], [])
                for para in tf.paragraphs:
                    for run in para.runs:
                        bold_vals.append(run.font.bold)
                        size_vals.append(run.font.size)
                        underline_vals.append(run.font.underline)
                        italic_vals.append(run.font.italic)
                if bold_vals:
                    size_emu = size_vals[0]
                    size_pt = round(size_emu / 12700, 1) if size_emu else None
                    return {'bold': bold_vals[0], 'size_pt': size_pt, 'underline': underline_vals[0], 'italic': italic_vals[0], 'text': tf.text.strip()[:60]}
        return {'error': 'No text shapes found on slide 1'}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_strikethrough_state__7f3881fcd58bcba327d08f20ed190dcc_qw35sft2_a43cff2a(env, config: dict):
    """Get strikethrough state for Finance Meetings and Program Coordinator bullets on slide 5."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file('/home/user/Desktop/New_Club_Spring_2018_Training.pptx')
        if not file_bytes:
            return {'error': 'file not found'}
        prs = Presentation(io.BytesIO(file_bytes))
        slide = prs.slides[4]
        state = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                para_text = para.text.strip()
                if not para_text:
                    continue
                runs = [r for r in para.runs if r.text.strip()]
                if runs:
                    has_strike = all((r.font.strike is True for r in runs))
                    state[para_text] = has_strike
        return state
    except Exception as e:
        return {'error': str(e)}

def get_writer_footer_body_spacing__e86d756354df667eec5758969f058362_qw35sft2_2e25e64c(env, config: dict):
    """
    Get footer page-number presence and line spacing of the paragraph
    starting with 'As an open-source replacement' (first real body paragraph).
    """
    path = config.get('path', '/home/user/Desktop/LibreOffice_Open_Source_Word_Processing.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_field': False, 'body_spacing_1_5': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_page_field = False
        with zipfile.ZipFile(tmp_path) as z:
            footer_files = [n for n in z.namelist() if re.match('word/footer\\d*\\.xml', n)]
            for fname in footer_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                if re.search('\\bPAGE\\b', content):
                    has_page_field = True
                    break
        from docx import Document
        doc = Document(tmp_path)
        body_spacing_1_5 = False
        for para in doc.paragraphs:
            if para.text.strip().startswith('As an open-source replacement'):
                para_xml = para._element.xml
                if re.search('w:line="360"', para_xml) and re.search('w:lineRule="auto"', para_xml):
                    body_spacing_1_5 = True
                try:
                    from docx.enum.text import WD_LINE_SPACING
                    rule = para.paragraph_format.line_spacing_rule
                    if rule == WD_LINE_SPACING.ONE_POINT_FIVE:
                        body_spacing_1_5 = True
                except Exception:
                    pass
                break
        return {'has_page_field': has_page_field, 'body_spacing_1_5': body_spacing_1_5}
    finally:
        os.unlink(tmp_path)

def get_title_run_font__ac8fd7a53908049060d14f1a620d5f79_qw35sft2_4c453cf5(env, config: dict):
    """Get the font name used in the first run of the first non-empty paragraph (title)."""
    import tempfile, os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/The Wonders of Our Solar System.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'font_name': None}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for para in doc.paragraphs:
            if para.text.strip() and para.runs:
                run = para.runs[0]
                font_name = run.font.name
                if font_name is None:
                    rPr = run._r.find(qn('w:rPr'))
                    if rPr is not None:
                        rFonts = rPr.find(qn('w:rFonts'))
                        if rFonts is not None:
                            font_name = rFonts.get(qn('w:ascii')) or rFonts.get(qn('w:hAnsi'))
                return {'font_name': font_name, 'text_preview': para.text[:60]}
        return {'font_name': None}
    except Exception as e:
        return {'error': str(e), 'font_name': None}
    finally:
        os.unlink(tmp_path)

def get_para0_underline__0ee57f85e0ff64364a5561938bc94f89_qw35sft2_eeee8dff(env, config: dict):
    """Get underline state of all runs in the first paragraph of the tutorial guidelines doc."""
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/CCCH9003_Tutorial_guidelines.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        if not para.runs:
            return {'underline': False, 'run_count': 0}
        all_underline = all((run.underline is True for run in para.runs))
        return {'underline': all_underline, 'run_count': len(para.runs)}
    finally:
        os.unlink(tmp_path)

def get_docx_last_line__9ae69d0ef002958997c53c817727d757_qw35sft2_1cd0991b(env, config: dict):
    """Get the last non-empty line of a docx file from the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if not lines:
            return {'last_line': '', 'total_lines': 0}
        return {'last_line': lines[-1], 'total_lines': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_docx_subscript_and_title_italic__5faa01d8a062ee366a1e56b39b81f8d9_qw35sft2_8cf73d1b(env, config: dict):
    """
    Get formatting state of H2O_Factsheet_WA.docx:
    - Whether the '2' in 'H2O' in the title paragraph has subscript formatting
    - Whether the '—Soak up the Science' portion of the title is italic
    """
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/H2O_Factsheet_WA.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        subscript_2_in_title = False
        title_soak_italic = False
        for para in doc.paragraphs:
            if para.style.name == 'Title':
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_2_in_title = True
                    if 'Soak' in run.text and run.italic is True:
                        title_soak_italic = True
                break
        return {'subscript_2_in_title': subscript_2_in_title, 'title_soak_italic': title_soak_italic}
    finally:
        os.unlink(tmp_path)

def get_writer_basic_fonts__cb1a08b69b2e8e76cab6534a166b3ea6_qw35sft2_3329d88b(env, config: dict):
    """
    Read LibreOffice XCU config and return the Standard (Default) and
    Heading font names from Basic Fonts (Western) settings.
    """
    result = {'default_font': None, 'heading_font': None}
    xcu_bytes = env.controller.get_file('/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    if not xcu_bytes:
        return result
    with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False) as f:
        f.write(xcu_bytes)
        xcu_path = f.name
    try:
        tree = ET.parse(xcu_path)
        root = tree.getroot()
        ns = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', ns):
            for prop in elem.findall('.//prop[@oor:name="Standard"]', ns):
                for value in prop.findall('value', ns):
                    result['default_font'] = value.text
            for prop in elem.findall('.//prop[@oor:name="Heading"]', ns):
                for value in prop.findall('value', ns):
                    result['heading_font'] = value.text
    except Exception:
        pass
    finally:
        os.unlink(xcu_path)
    return result

def get_docx_lower_and_page_nums__4155f473baa624a7c08f0169367bee57_qw35sft2_721e17f7(env, config: dict):
    """Download docx and check lowercase text plus presence of page numbers in footer."""
    import tempfile
    import os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        alpha_chars = []
        for para in doc.paragraphs:
            for c in para.text:
                if c.isalpha():
                    alpha_chars.append(c)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for c in cell.text:
                        if c.isalpha():
                            alpha_chars.append(c)
        all_lower = all((c.islower() for c in alpha_chars)) if alpha_chars else False
        has_page_numbers = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            try:
                footer_xml = footer._element.xml
                if 'PAGE' in footer_xml:
                    has_page_numbers = True
                    break
            except Exception:
                pass
        return {'all_lower': all_lower, 'has_page_numbers': has_page_numbers, 'total_alpha': len(alpha_chars)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_writer_titlecase_center__c181651f58b43b69af608583a0523891_qw35sft2_1b32c475(env, config: dict):
    """Download Geography_And_Magical_Realism.docx and check:
    1. Title case applied across all content paragraphs.
    2. Center alignment on the title paragraph (para 0).
    """
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Geography_And_Magical_Realism.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        title_case_ok = True
        for para in paras:
            text = para.text.strip()
            if not text:
                continue
            for w in text.split():
                if w and w[0].isalpha() and (not w[0].isupper()):
                    title_case_ok = False
                    break
            if not title_case_ok:
                break
        title_para = paras[0]
        alignment = title_para.alignment
        title_centered = alignment is not None and int(alignment) == 1
        return {'title_case_applied': title_case_ok, 'title_centered': title_centered}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_align__2b822025bba1f4240036b27d0339b871_qw35sft2_ac7f2755(env, config: dict):
    """Read docx from VM and extract font name and paragraph alignment stats."""
    import tempfile
    import os
    from collections import Counter
    path = config.get('path', '/home/user/Desktop/Dublin_Zoo_Intro.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        doc = Document(tmp_path)
        font_names = []
        non_empty_paras = [p for p in doc.paragraphs if p.text.strip()]
        for para in non_empty_paras:
            for run in para.runs:
                if run.text.strip() and run.font.name:
                    font_names.append(run.font.name)
        most_common_font = Counter(font_names).most_common(1)[0][0] if font_names else None
        all_same_font = len(set(font_names)) <= 1 if font_names else False
        if non_empty_paras:
            center_count = sum((1 for p in non_empty_paras if p.alignment == WD_ALIGN_PARAGRAPH.CENTER))
            all_center = center_count >= len(non_empty_paras) * 0.9
        else:
            all_center = False
        return {'font_name': most_common_font, 'all_same_font': all_same_font, 'all_center_aligned': all_center}
    finally:
        os.unlink(tmp_path)

def get_writer_heading_body_align__d846f8a9d2cb6584815a6f17cfb8e80a_qw35sft2_0cdf8986(env, config: dict):
    """Get alignment of heading (para 0) and body text (para 2) from the Constitution docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Constitution_Template_With_Guidelines.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) < 3:
            return {'error': 'Not enough paragraphs'}
        heading_para = doc.paragraphs[0]
        body_para = doc.paragraphs[2]
        heading_centered = heading_para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        body_centered = body_para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        return {'heading_centered': heading_centered, 'body_centered': body_centered}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_first_para_strike__efe155f44671b1c09a26d5fefbd2fc44_qw35sft2_f933e3da(env, config: dict):
    import tempfile, os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/GEOG2169_Course_Outline_2022-23.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_strikethrough': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content_paras = [p for p in doc.paragraphs if p.text.strip()]
        if not content_paras:
            return {'error': 'No content paragraphs found', 'has_strikethrough': False}
        first_para = content_paras[0]
        runs_with_text = [r for r in first_para.runs if r.text.strip()]
        if not runs_with_text:
            return {'has_strikethrough': False, 'text': first_para.text[:100]}
        all_strike = all((bool(r.font.strike) for r in runs_with_text))
        return {'has_strikethrough': all_strike, 'text': first_para.text[:100]}
    finally:
        os.unlink(tmp_path)

def get_docx_italic_font_size__b8546f99d1a88f06fb0b6ffbfa13cf55_qw35sft2_e5a0021c(env, config: dict):
    """Download Y22-2119-assign4.docx and check whether all italic runs are 14pt."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Y22-2119-assign4.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_sizes = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic and run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    italic_sizes.append(size_pt)
        italic_count = len(italic_sizes)
        all_14 = all((s == 14.0 for s in italic_sizes)) if italic_sizes else False
        return {'italic_count': italic_count, 'all_italic_size_14': all_14}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_word_font_colors__336f1c4245e680e1aace133fb243059a_qw35sft2_3479f659(env, config: dict):
    """Extract font colors for all words in the table of the docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        word_colors = {}
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        text = para.text.strip().lower()
                        if not text:
                            continue
                        color_hex = None
                        for run in para.runs:
                            try:
                                if run.font.color.rgb is not None:
                                    color_hex = str(run.font.color.rgb).upper()
                            except Exception:
                                pass
                            break
                        word_colors[text] = color_hex
        return {'word_colors': word_colors}
    except Exception as e:
        return {'error': str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_odt_highlight_italic__e4e69693f939d9cc32e39d11fb21f92a_qw35sft2_45bc9ddb(env, config: dict):
    path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Could not load ODT file'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as zf:
            content = zf.read('content.xml').decode('utf-8')
            styles_xml = zf.read('styles.xml').decode('utf-8')
    finally:
        os.unlink(tmp_path)
    NS_S = '{' + _NS_STYLE_qw35sft2_7af95d + '}'
    NS_FO_ATTR = '{' + _NS_FO_qw35sft2_7af95d + '}'
    NS_T = '{' + _NS_TEXT_qw35sft2_7af95d + '}'
    NS_LOEXT_ATTR = '{' + _NS_LOEXT_qw35sft2_7af95d + '}'
    root_content = None
    root_styles = None
    try:
        root_content = ET.fromstring(content)
    except ET.ParseError:
        pass
    try:
        root_styles = ET.fromstring(styles_xml)
    except ET.ParseError:
        pass
    has_highlights = False
    for root_node in [root_content, root_styles]:
        if root_node is None or has_highlights:
            break
        for tp in root_node.iter(NS_S + 'text-properties'):
            loext_hl = tp.get(NS_LOEXT_ATTR + 'char-highlight-color', '')
            if loext_hl and loext_hl.lower() not in _HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d:
                has_highlights = True
                break
            bg = tp.get(NS_FO_ATTR + 'background-color', '')
            if bg and bg.lower() not in _HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d:
                has_highlights = True
                break
    title_text = 'Sample Recruitment Phone Script'
    title_para_node = None
    title_style_names = []
    if root_content is not None:
        for para in root_content.iter(NS_T + 'p'):
            all_text = ''.join(para.itertext())
            if title_text in all_text:
                title_para_node = para
                para_style = para.get(NS_T + 'style-name', '')
                if para_style:
                    title_style_names.append(para_style)
                for span in para.iter(NS_T + 'span'):
                    span_style = span.get(NS_T + 'style-name', '')
                    if span_style:
                        title_style_names.append(span_style)
                break
    all_styles = {}
    for root_node in [root_content, root_styles]:
        if root_node is None:
            continue
        for style in root_node.iter(NS_S + 'style'):
            sname = style.get(NS_S + 'name')
            if sname:
                all_styles[sname] = style
    title_italic = False
    if title_para_node is not None and (not title_italic):
        for tp in title_para_node.iter(NS_S + 'text-properties'):
            if tp.get(NS_FO_ATTR + 'font-style') == 'italic':
                title_italic = True
                break
    if not title_italic:
        for sname in title_style_names:
            visited = set()
            current = sname
            while current and current not in visited:
                visited.add(current)
                style_elem = all_styles.get(current)
                if style_elem is None:
                    break
                for tp in style_elem.iter(NS_S + 'text-properties'):
                    if tp.get(NS_FO_ATTR + 'font-style') == 'italic':
                        title_italic = True
                        break
                if title_italic:
                    break
                current = style_elem.get(NS_S + 'parent-style-name', '')
            if title_italic:
                break
    return {'has_highlights': has_highlights, 'title_italic': title_italic}

def get_docx_three_para_spacing__e4c36138929fcbc781e762a8996148c1_qw35sft2_3d322a41(env, config: dict):
    """Download Novels_Intro_Packet.docx and return line spacing for first three content paragraphs."""
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        result = {}
        for key, idx in (('para0_spacing', 0), ('para2_spacing', 2), ('para4_spacing', 4)):
            para = paras[idx]
            pPr = para._p.pPr
            if pPr is None:
                result[key] = 'single'
                continue
            sp = pPr.find(qn('w:spacing'))
            if sp is None:
                result[key] = 'single'
                continue
            line = sp.get(qn('w:line'))
            line_rule = sp.get(qn('w:lineRule'))
            if line_rule in ('auto', None):
                if line == '480':
                    result[key] = 'double'
                elif line == '360':
                    result[key] = '1.5'
                elif line == '240':
                    result[key] = 'single'
                else:
                    result[key] = 'other'
            else:
                result[key] = 'other'
        return result
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_extended_state__d40d48ebd32c5cdd693a8c4ab4565d6c_qw35sft2_1970bdc2(env, config: dict):
    """Read docx and check Steinberg presence, add-here removal, and document main title."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = doc.paragraphs
        full_text = '\n'.join((p.text for p in paragraphs))
        has_steinberg = 'Steinberg' in full_text
        has_add_here = '<add here>' in full_text
        main_title = ''
        for p in paragraphs:
            if p.style.name == 'Heading 1' and p.text.strip():
                main_title = p.text.strip()
                break
        return {'has_steinberg': has_steinberg, 'has_add_here_marker': has_add_here, 'main_title': main_title}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_spacing_arial__d472f99c5ab0bbd719c99178119f321f_qw35sft2_85781c1c(env, config: dict):
    """Get line spacings for 3 paragraphs and font names used in runs."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/CCHU9045_Course_Outline_2019-20.docx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        if len(paragraphs) < 3:
            return {'error': f'Expected at least 3 non-empty paragraphs, found {len(paragraphs)}'}
        intro_spacing = paragraphs[0].paragraph_format.line_spacing
        body_spacing = paragraphs[1].paragraph_format.line_spacing
        conclusion_spacing = paragraphs[2].paragraph_format.line_spacing
        all_font_names = []
        for para in paragraphs[:3]:
            for run in para.runs:
                if run.font.name is not None:
                    all_font_names.append(run.font.name)
        all_arial = len(all_font_names) > 0 and all((name == 'Arial' for name in all_font_names))
        return {'intro_spacing': float(intro_spacing) if intro_spacing is not None else None, 'body_spacing': float(body_spacing) if body_spacing is not None else None, 'conclusion_spacing': float(conclusion_spacing) if conclusion_spacing is not None else None, 'all_arial': all_arial}
    finally:
        os.unlink(tmp_path)

def get_writer_break_and_notes__4e9ea3f60c13559abc5abe689f4b9ea2_qw35sft2_57942942(env, config: dict):
    """
    Check Sample_Statutory_Declaration.docx for:
    1. Number of explicit run-level page breaks.
    2. Whether any paragraph text (after the 4th page break) contains 'Notes'.
    """
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Sample_Statutory_Declaration.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        explicit_breaks = 0
        break_indices = []
        paragraphs = doc.paragraphs
        for idx, para in enumerate(paragraphs):
            for run in para.runs:
                for br in run._element.findall('.//' + qn('w:br')):
                    if br.get(qn('w:type')) == 'page':
                        explicit_breaks += 1
                        break_indices.append(idx)
        notes_found = False
        if explicit_breaks >= 5:
            for bi in break_indices:
                for j in range(bi + 1, min(bi + 5, len(paragraphs))):
                    if 'notes' in paragraphs[j].text.lower():
                        notes_found = True
                        break
                if notes_found:
                    break
        return {'explicit_page_breaks': explicit_breaks, 'notes_found': notes_found}
    finally:
        os.unlink(tmp_path)

def get_writer_default_and_list_fonts__12134317917b6b593b5731418bb79a41_qw35sft2_75bbc093(env, config: dict):
    """
    Read LibreOffice XCU config and return the Standard (Default) and
    List font names from Basic Fonts (Western) settings.
    """
    result = {'default_font': None, 'list_font': None}
    xcu_bytes = env.controller.get_file('/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    if not xcu_bytes:
        return result
    with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False) as f:
        f.write(xcu_bytes)
        xcu_path = f.name
    try:
        tree = ET.parse(xcu_path)
        root = tree.getroot()
        ns = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', ns):
            for prop in elem.findall('.//prop[@oor:name="Standard"]', ns):
                for value in prop.findall('value', ns):
                    result['default_font'] = value.text
            for prop in elem.findall('.//prop[@oor:name="List"]', ns):
                for value in prop.findall('value', ns):
                    result['list_font'] = value.text
    except Exception:
        pass
    finally:
        os.unlink(xcu_path)
    return result

def get_docx_first_line__d29d0e1e9fda0acc872853f3a63d2906_qw35sft2_69a56e4c(env, config: dict):
    """Get the first non-empty line of a docx file from the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if not lines:
            return {'first_line': '', 'total_lines': 0}
        return {'first_line': lines[0], 'total_lines': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_writer_footer_title_align__54873b40ed0174b0628f230c8ad47868_qw35sft2_ecff64b7(env, config: dict):
    """Get footer page-number presence and title paragraph alignment."""
    path = config.get('path', '/home/user/Desktop/LibreOffice_Open_Source_Word_Processing.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_field': False, 'title_centered': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_page_field = False
        with zipfile.ZipFile(tmp_path) as z:
            footer_files = [n for n in z.namelist() if re.match('word/footer\\d*\\.xml', n)]
            for fname in footer_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                if re.search('\\bPAGE\\b', content):
                    has_page_field = True
                    break
        from docx import Document
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        doc = Document(tmp_path)
        title_para = doc.paragraphs[0]
        para_xml = title_para._element.xml
        title_centered = bool(re.search('w:jc\\s+w:val="center"', para_xml, re.IGNORECASE))
        return {'has_page_field': has_page_field, 'title_centered': title_centered}
    finally:
        os.unlink(tmp_path)

def get_docx_lower_and_bold_title__5844d45da47ec52046a23219bd5db770_qw35sft2_ab4c298e(env, config: dict):
    """Download docx, check all text is lowercase and the first non-empty
    paragraph's runs are all bold."""
    import tempfile
    import os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        alpha_chars = []
        for para in doc.paragraphs:
            for c in para.text:
                if c.isalpha():
                    alpha_chars.append(c)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for c in cell.text:
                        if c.isalpha():
                            alpha_chars.append(c)
        all_lower = all((c.islower() for c in alpha_chars)) if alpha_chars else False
        first_para_bold = False
        first_para_text = None
        for para in doc.paragraphs:
            if para.text.strip():
                first_para_text = para.text.strip()
                runs = para.runs
                if runs:
                    first_para_bold = all((run.bold is True for run in runs if run.text.strip()))
                else:
                    first_para_bold = False
                break
        return {'all_lower': all_lower, 'first_para_bold': first_para_bold, 'first_para_text': first_para_text, 'total_alpha': len(alpha_chars)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_docx_subscript_and_title_center__49f57bac6148724412fcfb3eb162539d_qw35sft2_4fe9563d(env, config: dict):
    """
    Get formatting state of H2O_Factsheet_WA.docx:
    - Whether the '2' in 'H2O' in the title paragraph has subscript formatting
    - Whether the title paragraph is center-aligned (alignment == 1)
    """
    import tempfile
    import os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_bytes = env.controller.get_file('/home/user/Desktop/H2O_Factsheet_WA.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        subscript_2_in_title = False
        title_centered = False
        for para in doc.paragraphs:
            if para.style.name == 'Title':
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_2_in_title = True
                        break
                title_centered = para.alignment == WD_ALIGN_PARAGRAPH.CENTER
                break
        return {'subscript_2_in_title': subscript_2_in_title, 'title_centered': title_centered}
    finally:
        os.unlink(tmp_path)

def get_writer_titlecase_italic__7a4ad09a1377c8f0eb798dbce87a7dce_qw35sft2_1af9ffcd(env, config: dict):
    """Download Geography_And_Magical_Realism.docx and check:
    1. Title case applied across all content paragraphs.
    2. Italic formatting on the title paragraph (para 0) runs.
    """
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Geography_And_Magical_Realism.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        title_case_ok = True
        for para in paras:
            text = para.text.strip()
            if not text:
                continue
            for w in text.split():
                if w and w[0].isalpha() and (not w[0].isupper()):
                    title_case_ok = False
                    break
            if not title_case_ok:
                break
        title_para = paras[0]
        title_italic = False
        if title_para.runs:
            title_italic = all((run.italic is True for run in title_para.runs))
        return {'title_case_applied': title_case_ok, 'title_italic': title_italic}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_multi_para_strike__d8440cb16db3c65fff5a2b530ef73072_qw35sft2_8f67219f(env, config: dict):
    import tempfile, os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/GEOG2169_Course_Outline_2022-23.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'first_para_strike': False, 'last_para_strike': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content_paras = [p for p in doc.paragraphs if p.text.strip()]
        if len(content_paras) < 2:
            return {'error': 'Not enough paragraphs', 'first_para_strike': False, 'last_para_strike': False}
        first_para = content_paras[0]
        last_para = content_paras[-1]
        first_runs = [r for r in first_para.runs if r.text.strip()]
        last_runs = [r for r in last_para.runs if r.text.strip()]
        first_strike = bool(first_runs) and all((bool(r.font.strike) for r in first_runs))
        last_strike = bool(last_runs) and all((bool(r.font.strike) for r in last_runs))
        return {'first_para_strike': first_strike, 'last_para_strike': last_strike, 'first_para_text': first_para.text[:80], 'last_para_text': last_para.text[:80]}
    finally:
        os.unlink(tmp_path)

def get_docx_italic_size_underline__dce073995a5927fc743181d7c02c0659_qw35sft2_3d981ede(env, config: dict):
    """Download Y22-2119-assign4.docx and check whether all italic runs are 14pt and underlined."""
    import tempfile
    import os
    from docx import Document
    from docx.enum.text import WD_UNDERLINE
    path = config.get('path', '/home/user/Desktop/Y22-2119-assign4.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_runs = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic and run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    underline = run.font.underline
                    is_underlined = bool(underline) if underline is not None else False
                    italic_runs.append({'size': size_pt, 'underlined': is_underlined})
        if not italic_runs:
            return {'italic_count': 0, 'all_italic_size_14': False, 'all_italic_underlined': False}
        all_14 = all((r['size'] == 14.0 for r in italic_runs))
        all_underlined = all((r['underlined'] for r in italic_runs))
        return {'italic_count': len(italic_runs), 'all_italic_size_14': all_14, 'all_italic_underlined': all_underlined}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_line_spacing__c985596bc3953d61c7d274fc30ba4960_qw35sft2_14feb6d7(env, config: dict):
    """Download Novels_Intro_Packet.docx and return line spacing for first two content paragraphs."""
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        results = {}
        for key, para in (('para0_spacing', paras[0]), ('para2_spacing', paras[2])):
            pPr = para._p.pPr
            if pPr is None:
                results[key] = 'single'
                continue
            sp = pPr.find(qn('w:spacing'))
            if sp is None:
                results[key] = 'single'
                continue
            line = sp.get(qn('w:line'))
            line_rule = sp.get(qn('w:lineRule'))
            if line_rule in ('auto', None):
                if line == '480':
                    results[key] = 'double'
                elif line == '360':
                    results[key] = '1.5'
                elif line == '240':
                    results[key] = 'single'
                else:
                    results[key] = 'other'
            else:
                results[key] = 'other'
        return results
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_size__a567f4bdee19e51d5ed6a587d8d38f9d_qw35sft2_ae1290b0(env, config: dict):
    """Read docx from VM and extract font name and font size stats across all runs."""
    import tempfile
    import os
    from collections import Counter
    path = config.get('path', '/home/user/Desktop/Dublin_Zoo_Intro.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        font_names = []
        font_sizes = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    if run.font.name:
                        font_names.append(run.font.name)
                    if run.font.size:
                        font_sizes.append(round(run.font.size.pt, 1))
        most_common_font = Counter(font_names).most_common(1)[0][0] if font_names else None
        most_common_size = Counter(font_sizes).most_common(1)[0][0] if font_sizes else None
        all_same_font = len(set(font_names)) <= 1 if font_names else False
        all_same_size = len(set(font_sizes)) <= 1 if font_sizes else False
        return {'font_name': most_common_font, 'font_size_pt': most_common_size, 'all_same_font': all_same_font, 'all_same_size': all_same_size}
    finally:
        os.unlink(tmp_path)

def get_writer_heading_align_size__a3c76448cfe24f131aa7bba8054b5d03_qw35sft2_7bb7f6c5(env, config: dict):
    """Get heading alignment and explicit font size from the Constitution docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Constitution_Template_With_Guidelines.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        is_centered = para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        font_size_pt = None
        for run in para.runs:
            if run.font.size is not None:
                font_size_pt = run.font.size / 12700.0
                break
        return {'is_centered': is_centered, 'font_size_pt': font_size_pt}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_word_font_colors__d89e455c3ca4dbc17e773411cf8f66df_qw35sft2_0501c7f4(env, config: dict):
    """Extract font colors for all words in the table of the docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        word_colors = {}
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        text = para.text.strip().lower()
                        if not text:
                            continue
                        color_hex = None
                        for run in para.runs:
                            try:
                                if run.font.color.rgb is not None:
                                    color_hex = str(run.font.color.rgb).upper()
                            except Exception:
                                pass
                            break
                        word_colors[text] = color_hex
        return {'word_colors': word_colors}
    except Exception as e:
        return {'error': str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_odt_highlight_font__582984f0118e076c9ea08b5f36393d3b_qw35sft2_afca0b7a(env, config: dict):
    path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Could not load ODT file'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as zf:
            content = zf.read('content.xml').decode('utf-8')
            styles_xml = zf.read('styles.xml').decode('utf-8')
    finally:
        os.unlink(tmp_path)
    NS_S = '{urn:oasis:names:tc:opendocument:xmlns:style:1.0}'
    NS_FO = '{urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0}'
    NS_LOEXT = '{urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0}'
    NS_TEXT = '{urn:oasis:names:tc:opendocument:xmlns:text:1.0}'
    no_hl_values = {'', 'transparent', 'automatic', '#00000000', 'none', '-1', '4294967295'}
    has_highlights = False
    try:
        root = ET.fromstring(content)
        for tp in root.iter(NS_S + 'text-properties'):
            bg = (tp.get(NS_FO + 'background-color') or '').strip().lower()
            loext_hl = (tp.get(NS_LOEXT + 'char-highlight-color') or '').strip().lower()
            if bg and bg not in no_hl_values or (loext_hl and loext_hl not in no_hl_values):
                has_highlights = True
                break
        if not has_highlights:
            for span in root.iter(NS_TEXT + 'span'):
                loext_hl = (span.get(NS_LOEXT + 'char-highlight-color') or '').strip().lower()
                bg = (span.get(NS_FO + 'background-color') or '').strip().lower()
                if bg and bg not in no_hl_values or (loext_hl and loext_hl not in no_hl_values):
                    has_highlights = True
                    break
    except ET.ParseError:
        pass
    default_font = None
    try:
        styles_root = ET.fromstring(styles_xml)
        for ds in styles_root.iter(NS_S + 'default-style'):
            if ds.get(NS_S + 'family') == 'paragraph':
                for tp in ds.iter(NS_S + 'text-properties'):
                    font = tp.get(NS_S + 'font-name') or tp.get(NS_FO + 'font-family')
                    if font:
                        default_font = font
                        break
            if default_font:
                break
        if default_font is None:
            for style in styles_root.iter(NS_S + 'style'):
                if style.get(NS_S + 'name') == 'Default Paragraph Style':
                    for tp in style.iter(NS_S + 'text-properties'):
                        font = tp.get(NS_S + 'font-name') or tp.get(NS_FO + 'font-family')
                        if font:
                            default_font = font
                            break
                if default_font:
                    break
    except ET.ParseError:
        pass
    if default_font is None:
        try:
            content_root = ET.fromstring(content)
            font_counts = {}
            for style in content_root.iter(NS_S + 'style'):
                for tp in style.iter(NS_S + 'text-properties'):
                    font = tp.get(NS_S + 'font-name') or tp.get(NS_FO + 'font-family')
                    if font:
                        font_counts[font] = font_counts.get(font, 0) + 1
            if font_counts:
                default_font = max(font_counts, key=font_counts.get)
        except ET.ParseError:
            pass
    return {'has_highlights': has_highlights, 'default_font': default_font}

def get_writer_three_state__c8e0b7dd7367f091fd322e812e9986d0_qw35sft2_98c08613(env, config: dict):
    """Read docx and return Steinberg presence, add-here removal, and References heading text."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join((p.text for p in doc.paragraphs))
        has_steinberg = 'Steinberg' in full_text
        has_add_here = '<add here>' in full_text
        refs_heading_text = ''
        for p in doc.paragraphs:
            if p.style.name.startswith('Heading') and ('Reference' in p.text or 'Bibliography' in p.text):
                refs_heading_text = p.text.strip()
                break
        return {'has_steinberg': has_steinberg, 'has_add_here_marker': has_add_here, 'refs_heading': refs_heading_text}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_spacing_fontsize__122cebb30dd4bbe8d6067021381568cf_qw35sft2_355b80cf(env, config: dict):
    """Get line spacings and font sizes for the 3-paragraph essay document."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/CCHU9045_Course_Outline_2019-20.docx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        if len(paragraphs) < 3:
            return {'error': f'Expected at least 3 non-empty paragraphs, found {len(paragraphs)}'}
        intro_spacing = paragraphs[0].paragraph_format.line_spacing
        body_spacing = paragraphs[1].paragraph_format.line_spacing
        conclusion_spacing = paragraphs[2].paragraph_format.line_spacing
        all_run_sizes = []
        for para in paragraphs[:3]:
            for run in para.runs:
                if run.font.size is not None:
                    all_run_sizes.append(int(run.font.size))
        all_14pt = len(all_run_sizes) > 0 and all((sz == 177800 for sz in all_run_sizes))
        return {'intro_spacing': float(intro_spacing) if intro_spacing is not None else None, 'body_spacing': float(body_spacing) if body_spacing is not None else None, 'conclusion_spacing': float(conclusion_spacing) if conclusion_spacing is not None else None, 'all_14pt': all_14pt}
    finally:
        os.unlink(tmp_path)

def get_writer_break_and_font__52445c32d012b75894f7e753d9bb73ed_qw35sft2_9ef5f0c7(env, config: dict):
    """
    Check Sample_Statutory_Declaration.docx for:
    1. Number of explicit run-level page breaks.
    2. Font name of the title paragraph's first run.
    """
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Sample_Statutory_Declaration.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        explicit_breaks = 0
        for para in doc.paragraphs:
            for run in para.runs:
                for br in run._element.findall('.//' + qn('w:br')):
                    if br.get(qn('w:type')) == 'page':
                        explicit_breaks += 1
        title_font = None
        title_para = doc.paragraphs[0]
        if title_para.runs:
            title_font = title_para.runs[0].font.name
        return {'explicit_page_breaks': explicit_breaks, 'title_font': title_font}
    finally:
        os.unlink(tmp_path)

def get_title_alignment__6435e9b69d0bbcb863c89964e3084688_qw35sft2_568a8a6b(env, config: dict):
    """Get the paragraph alignment of the first non-empty paragraph (title) in the document."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    path = config.get('path', '/home/user/Desktop/The Wonders of Our Solar System.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'alignment': None}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    alignment_names = {0: 'LEFT', 1: 'CENTER', 2: 'RIGHT', 3: 'JUSTIFY', None: 'LEFT'}
    try:
        doc = Document(tmp_path)
        for para in doc.paragraphs:
            if para.text.strip():
                align_val = para.alignment
                if align_val is None:
                    alignment = 'LEFT'
                else:
                    alignment = alignment_names.get(int(align_val), str(align_val))
                return {'alignment': alignment, 'text_preview': para.text[:60]}
        return {'alignment': 'LEFT', 'text_preview': ''}
    except Exception as e:
        return {'error': str(e), 'alignment': None}
    finally:
        os.unlink(tmp_path)

def get_writer_titlecase_underline__d22455594f3cb66658bea4571465f8f1_qw35sft2_d8cd9c76(env, config: dict):
    """Download Geography_And_Magical_Realism.docx and check:
    1. Title case applied across all content paragraphs.
    2. Underline formatting on the title paragraph (para 0) runs.
    """
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Geography_And_Magical_Realism.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        title_case_ok = True
        for para in paras:
            text = para.text.strip()
            if not text:
                continue
            for w in text.split():
                if w and w[0].isalpha() and (not w[0].isupper()):
                    title_case_ok = False
                    break
            if not title_case_ok:
                break
        title_para = paras[0]
        title_underline = False
        if title_para.runs:
            title_underline = all((run.font.underline is True for run in title_para.runs))
        return {'title_case_applied': title_case_ok, 'title_underline': title_underline}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_footer_pageno__035b6f0ae922d6e2d1f24326ec60c904_qw35sft2_865faa47(env, config: dict):
    """Get footer page-number field presence from the Writer document."""
    path = config.get('path', '/home/user/Desktop/LibreOffice_Open_Source_Word_Processing.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_field': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_page_field = False
        with zipfile.ZipFile(tmp_path) as z:
            footer_files = [n for n in z.namelist() if re.match('word/footer\\d*\\.xml', n)]
            for fname in footer_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                if re.search('\\bPAGE\\b', content):
                    has_page_field = True
                    break
        return {'has_page_field': has_page_field}
    finally:
        os.unlink(tmp_path)

def get_docx_text_all_upper__cbc878ff0e7a736898047bffbd6f85f6_qw35sft2_5ca44771(env, config: dict):
    """Download docx and check whether all alphabetic characters are uppercase."""
    import tempfile
    import os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        alpha_chars = []
        for para in doc.paragraphs:
            for c in para.text:
                if c.isalpha():
                    alpha_chars.append(c)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for c in cell.text:
                        if c.isalpha():
                            alpha_chars.append(c)
        if not alpha_chars:
            return {'error': 'No alphabetic text found in document'}
        all_upper = all((c.isupper() for c in alpha_chars))
        upper_count = sum((1 for c in alpha_chars if c.isupper()))
        upper_ratio = upper_count / len(alpha_chars)
        return {'all_upper': all_upper, 'upper_ratio': upper_ratio, 'total_alpha': len(alpha_chars)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_para0_bold__42cf2947548fcbcc72e5e44a50eb60dd_qw35sft2_99cf7f57(env, config: dict):
    """Get bold state of all runs in the first paragraph of the tutorial guidelines doc."""
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/CCCH9003_Tutorial_guidelines.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        if not para.runs:
            return {'bold': False, 'run_count': 0}
        all_bold = all((run.bold is True for run in para.runs))
        return {'bold': all_bold, 'run_count': len(para.runs)}
    finally:
        os.unlink(tmp_path)

def get_docx_train_records__aa11b4f1fde2d6cc216fd8dac61371d5_qw35sft2_1468b0d6(env, config: dict):
    """Get train record stats from a docx file: presence/absence of a specific train ID and total line count."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        target_train = config.get('target_train', 'E229')
        target_count = sum((1 for l in lines if target_train in l.split(',')[1].strip() if len(l.split(',')) >= 2))
        return {'total_lines': len(lines), 'target_train': target_train, 'target_count': target_count}
    finally:
        os.unlink(tmp_path)

def get_writer_font_and_alignment__70647269807c892c9cd9454bb994336f_qw35sft2_5d687801(env, config: dict):
    """
    Fetch LibreOffice XCU config and the docx file, then return:
    - default_font: the Standard font name from XCU
    - first_para_centered: whether the first paragraph is center-aligned
    """
    result = {'default_font': None, 'first_para_centered': False}
    xcu_bytes = env.controller.get_file('/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    if xcu_bytes:
        with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False) as f:
            f.write(xcu_bytes)
            xcu_path = f.name
        try:
            tree = ET.parse(xcu_path)
            root = tree.getroot()
            ns = {'oor': 'http://openoffice.org/2001/registry'}
            for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', ns):
                for prop in elem.findall('.//prop[@oor:name="Standard"]', ns):
                    for value in prop.findall('value', ns):
                        result['default_font'] = value.text
        except Exception:
            pass
        finally:
            os.unlink(xcu_path)
    docx_bytes = env.controller.get_file('/home/user/Desktop/loa-one-time-submission-sealand.docx')
    if docx_bytes:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            f.write(docx_bytes)
            docx_path = f.name
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            doc = Document(docx_path)
            if doc.paragraphs:
                first_para = doc.paragraphs[0]
                alignment = first_para.paragraph_format.alignment
                result['first_para_centered'] = alignment == WD_ALIGN_PARAGRAPH.CENTER
        except Exception:
            pass
        finally:
            os.unlink(docx_path)
    return result

def get_docx_subscript_and_bold_heading__fa57f784af7171fa0e677f21369b4177_qw35sft2_c13dc097(env, config: dict):
    """
    Get formatting state of H2O_Factsheet_WA.docx:
    - Whether the '2' in 'H2O' in the title paragraph has subscript formatting
    - Whether the 'Fact sheet' heading paragraph has any bold run
    """
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/H2O_Factsheet_WA.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        subscript_2_in_title = False
        for para in doc.paragraphs:
            if para.style.name == 'Title':
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_2_in_title = True
                        break
                break
        bold_heading = False
        for para in doc.paragraphs:
            if para.style.name == 'Heading 1' and 'Fact sheet' in para.text:
                for run in para.runs:
                    if run.bold is True:
                        bold_heading = True
                        break
                break
        return {'subscript_2_in_title': subscript_2_in_title, 'bold_heading': bold_heading}
    finally:
        os.unlink(tmp_path)

def get_docx_italic_size_16__e2de398ff1a46d21a89c93f18cb16653_qw35sft2_5af2a69c(env, config: dict):
    """Download Y22-2119-assign4.docx and check whether all italic runs are 16pt."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Y22-2119-assign4.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_sizes = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic and run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    italic_sizes.append(size_pt)
        italic_count = len(italic_sizes)
        all_16 = all((s == 16.0 for s in italic_sizes)) if italic_sizes else False
        return {'italic_count': italic_count, 'all_italic_size_16': all_16}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_bold__3a9eb9ceb78a9dd18ebdb09636bdbe29_qw35sft2_786d155a(env, config: dict):
    """Read docx from VM and extract font name and bold state across all runs."""
    import tempfile
    import os
    from collections import Counter
    path = config.get('path', '/home/user/Desktop/Dublin_Zoo_Intro.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        font_names = []
        runs_with_text = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    runs_with_text.append(run)
                    if run.font.name:
                        font_names.append(run.font.name)
        most_common_font = Counter(font_names).most_common(1)[0][0] if font_names else None
        all_same_font = len(set(font_names)) <= 1 if font_names else False
        if runs_with_text:
            bold_count = sum((1 for r in runs_with_text if r.bold is True))
            all_bold = bold_count >= len(runs_with_text) * 0.9
        else:
            all_bold = False
        return {'font_name': most_common_font, 'all_same_font': all_same_font, 'all_bold': all_bold}
    finally:
        os.unlink(tmp_path)

def get_writer_heading_align_font__fb136ac47c87dcc7e787c09564b245bc_qw35sft2_7c722193(env, config: dict):
    """Get heading alignment and font name from the Constitution docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Constitution_Template_With_Guidelines.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        is_centered = para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        font_name = None
        for run in para.runs:
            if run.font.name:
                font_name = run.font.name
                break
        return {'is_centered': is_centered, 'font_name': font_name}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_strike_bold_chain__30608d866c349bd16efe53a1c72a27d7_qw35sft2_95f689f7(env, config: dict):
    import tempfile, os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/GEOG2169_Course_Outline_2022-23.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'last_para_strike': False, 'first_para_bold': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content_paras = [p for p in doc.paragraphs if p.text.strip()]
        if len(content_paras) < 2:
            return {'error': 'Not enough paragraphs', 'last_para_strike': False, 'first_para_bold': False}
        first_para = content_paras[0]
        last_para = content_paras[-1]
        first_runs = [r for r in first_para.runs if r.text.strip()]
        last_runs = [r for r in last_para.runs if r.text.strip()]
        first_bold = bool(first_runs) and all((bool(r.bold) for r in first_runs))
        last_strike = bool(last_runs) and all((bool(r.font.strike) for r in last_runs))
        return {'first_para_bold': first_bold, 'last_para_strike': last_strike, 'first_para_text': first_para.text[:80], 'last_para_text': last_para.text[:80]}
    finally:
        os.unlink(tmp_path)

def get_word_font_colors__e97799f0add4268fcc31cdd1e95ac277_qw35sft2_a5782795(env, config: dict):
    """Extract font colors for all words in the table of the docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        word_colors = {}
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        text = para.text.strip().lower()
                        if not text:
                            continue
                        color_hex = None
                        for run in para.runs:
                            try:
                                if run.font.color.rgb is not None:
                                    color_hex = str(run.font.color.rgb).upper()
                            except Exception:
                                pass
                            break
                        word_colors[text] = color_hex
        return {'word_colors': word_colors}
    except Exception as e:
        return {'error': str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_docx_mixed_spacing__ffe17a8bd50575f09eec01a68fed60c6_qw35sft2_6f5cace1(env, config: dict):
    """Download Novels_Intro_Packet.docx and return line spacing for first two content paragraphs."""
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Novels_Intro_Packet.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        results = {}
        for key, para in (('para0_spacing', paras[0]), ('para2_spacing', paras[2])):
            pPr = para._p.pPr
            if pPr is None:
                results[key] = 'single'
                continue
            sp = pPr.find(qn('w:spacing'))
            if sp is None:
                results[key] = 'single'
                continue
            line = sp.get(qn('w:line'))
            line_rule = sp.get(qn('w:lineRule'))
            if line_rule in ('auto', None):
                if line == '480':
                    results[key] = 'double'
                elif line == '360':
                    results[key] = '1.5'
                elif line == '240':
                    results[key] = 'single'
                else:
                    results[key] = 'other'
            else:
                results[key] = 'other'
        return results
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_odt_highlight_strikethrough__a30ceeb68a8eeb0b8a512ab61a2387e1_qw35sft2_8ca7650f(env, config: dict):
    path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Could not load ODT file'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as zf:
            content = zf.read('content.xml').decode('utf-8')
            styles_xml = zf.read('styles.xml').decode('utf-8')
    finally:
        os.unlink(tmp_path)
    NS_S = '{' + _NS_STYLE_qw35sft2_fd94cd + '}'
    NS_FO_B = '{' + _NS_FO_qw35sft2_fd94cd + '}'
    NS_T = '{' + _NS_TEXT_qw35sft2_fd94cd + '}'
    has_highlights = False
    try:
        hl_root = ET.fromstring(content)
        for style in hl_root.iter(NS_S + 'style'):
            if style.get(NS_S + 'family') != 'text':
                continue
            for tp in style.iter(NS_S + 'text-properties'):
                bg = tp.get(NS_FO_B + 'background-color', '')
                if bg and bg.lower() not in ('transparent', '', 'automatic', '#00000000'):
                    has_highlights = True
                    break
            if has_highlights:
                break
    except ET.ParseError:
        pass
    last_sentence = 'It was nice speaking with you'
    last_sentence_strikethrough = False
    style_names = []
    try:
        para_root = ET.fromstring(content)
        for para in para_root.iter(NS_T + 'p'):
            all_text = ''.join(para.itertext())
            if last_sentence in all_text:
                para_style = para.get(NS_T + 'style-name', '')
                if para_style:
                    style_names.append(para_style)
                for span in para.iter(NS_T + 'span'):
                    span_text = ''.join(span.itertext())
                    if last_sentence in span_text or len(span_text) > 5:
                        span_style = span.get(NS_T + 'style-name', '')
                        if span_style:
                            style_names.append(span_style)
                break
    except ET.ParseError:
        pass
    for sname in style_names:
        if last_sentence_strikethrough:
            break
        for xml_src in [content, styles_xml]:
            if last_sentence_strikethrough:
                break
            try:
                sroot = ET.fromstring(xml_src)
            except ET.ParseError:
                continue
            for style in sroot.iter(NS_S + 'style'):
                if style.get(NS_S + 'name') == sname:
                    for tp in style.iter(NS_S + 'text-properties'):
                        lt_style = dict(tp.attrib).get(NS_S + 'text-line-through-style', '')
                        if lt_style and lt_style not in ('none', ''):
                            last_sentence_strikethrough = True
                            break
                    break
    return {'has_highlights': has_highlights, 'last_sentence_strikethrough': last_sentence_strikethrough}

def get_writer_ref_state__dab123b7868d167252ed777095add0ba_qw35sft2_03ce8375(env, config: dict):
    """Read docx and check for Steinberg reference presence and <add here> marker removal."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join((p.text for p in doc.paragraphs))
        has_steinberg = 'Steinberg' in full_text
        has_add_here = '<add here>' in full_text
        return {'has_steinberg': has_steinberg, 'has_add_here_marker': has_add_here}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_spacing_bold_intro__7ab11a38694adeb4332550a69fd5abfe_qw35sft2_57061efb(env, config: dict):
    """Get line spacings for 3 paragraphs and bold status of introduction."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/CCHU9045_Course_Outline_2019-20.docx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        if len(paragraphs) < 3:
            return {'error': f'Expected at least 3 non-empty paragraphs, found {len(paragraphs)}'}
        intro_spacing = paragraphs[0].paragraph_format.line_spacing
        body_spacing = paragraphs[1].paragraph_format.line_spacing
        conclusion_spacing = paragraphs[2].paragraph_format.line_spacing
        intro_runs = paragraphs[0].runs
        intro_bold = len(intro_runs) > 0 and all((run.bold is True for run in intro_runs))
        return {'intro_spacing': float(intro_spacing) if intro_spacing is not None else None, 'body_spacing': float(body_spacing) if body_spacing is not None else None, 'conclusion_spacing': float(conclusion_spacing) if conclusion_spacing is not None else None, 'intro_bold': intro_bold}
    finally:
        os.unlink(tmp_path)

def get_writer_break_and_title_italic__0f7fdf2d11eca0ddff6ddca8ba1c7a92_qw35sft2_9f93c670(env, config: dict):
    """
    Check Sample_Statutory_Declaration.docx for:
    1. Number of explicit run-level page breaks.
    2. Whether the first paragraph (title) has italic formatting on its first run.
    """
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Sample_Statutory_Declaration.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        explicit_breaks = 0
        for para in doc.paragraphs:
            for run in para.runs:
                for br in run._element.findall('.//' + qn('w:br')):
                    if br.get(qn('w:type')) == 'page':
                        explicit_breaks += 1
        title_italic = False
        title_para = doc.paragraphs[0]
        if title_para.runs:
            run0 = title_para.runs[0]
            title_italic = bool(run0.italic)
        return {'explicit_page_breaks': explicit_breaks, 'title_italic': title_italic}
    finally:
        os.unlink(tmp_path)

def get_docx_sentence_case__121c5b95c7fa376e725e2a74f1f18a20_qw35sft2_0ffe6d01(env, config: dict):
    """Download docx and check if text follows sentence case (first alpha of each
    paragraph is uppercase, all other alpha characters are lowercase)."""
    import tempfile
    import os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        total_blocks = 0
        first_upper_ok = 0
        rest_lower_ok = 0
        texts_to_check = [para.text for para in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    texts_to_check.append(cell.text)
        for text in texts_to_check:
            alpha = [c for c in text if c.isalpha()]
            if not alpha:
                continue
            total_blocks += 1
            if alpha[0].isupper():
                first_upper_ok += 1
            if len(alpha) == 1 or all((c.islower() for c in alpha[1:])):
                rest_lower_ok += 1
        if total_blocks == 0:
            return {'error': 'No text blocks found'}
        return {'total_blocks': total_blocks, 'first_upper_ratio': first_upper_ok / total_blocks, 'rest_lower_ratio': rest_lower_ok / total_blocks, 'is_sentence_case': first_upper_ok == total_blocks and rest_lower_ok == total_blocks}
    except Exception as e:
        return {'error': str(e)}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_writer_footer_title_italic__6a6cc374d796e945ca1edd575879589a_qw35sft2_e128fbf1(env, config: dict):
    """Get footer page-number presence and whether the document title is italic."""
    path = config.get('path', '/home/user/Desktop/LibreOffice_Open_Source_Word_Processing.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_field': False, 'title_italic': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_page_field = False
        with zipfile.ZipFile(tmp_path) as z:
            footer_files = [n for n in z.namelist() if re.match('word/footer\\d*\\.xml', n)]
            for fname in footer_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                if re.search('\\bPAGE\\b', content):
                    has_page_field = True
                    break
        from docx import Document
        doc = Document(tmp_path)
        title_para = doc.paragraphs[0]
        title_italic = False
        for run in title_para.runs:
            run_xml = run._element.xml
            if re.search('<w:i(?!\\w)', run_xml):
                if not re.search('<w:i\\s+w:val="0"', run_xml):
                    title_italic = True
                    break
        return {'has_page_field': has_page_field, 'title_italic': title_italic}
    finally:
        os.unlink(tmp_path)

def get_docx_last_line__8efbbd8cc0ef91addd7e2f867e2ff2b5_qw35sft2_cefd8c72(env, config: dict):
    """Get the last non-empty line of a docx file from the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if not lines:
            return {'last_line': '', 'total_lines': 0}
        return {'last_line': lines[-1], 'total_lines': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_docx_last_line__e6176eae3a47bbdd8f12a06b8b082ed4_qw35sft2_cb11f3e2(env, config: dict):
    """Get the last non-empty line of a docx file from the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if not lines:
            return {'last_line': '', 'total_lines': 0}
        return {'last_line': lines[-1], 'total_lines': len(lines)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_and_pagebreak__e2b27ec72c919c5a4bd4f151e7ff64b6_qw35sft2_df0dbf5e(env, config: dict):
    """
    Fetch LibreOffice XCU config and the docx file, then return:
    - default_font: the Standard font name from XCU
    - page_break_count: number of explicit page breaks in the docx
    """
    result = {'default_font': None, 'page_break_count': 0}
    xcu_bytes = env.controller.get_file('/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    if xcu_bytes:
        with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False) as f:
            f.write(xcu_bytes)
            xcu_path = f.name
        try:
            tree = ET.parse(xcu_path)
            root = tree.getroot()
            ns = {'oor': 'http://openoffice.org/2001/registry'}
            for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', ns):
                for prop in elem.findall('.//prop[@oor:name="Standard"]', ns):
                    for value in prop.findall('value', ns):
                        result['default_font'] = value.text
        except Exception:
            pass
        finally:
            os.unlink(xcu_path)
    docx_bytes = env.controller.get_file('/home/user/Desktop/loa-one-time-submission-sealand.docx')
    if docx_bytes:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            f.write(docx_bytes)
            docx_path = f.name
        try:
            from docx import Document
            doc = Document(docx_path)
            namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            page_break_count = 0
            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    for br in run.element.findall('.//w:br', namespaces):
                        br_type = br.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type')
                        if br_type == 'page':
                            page_break_count += 1
            result['page_break_count'] = page_break_count
        except Exception:
            pass
        finally:
            os.unlink(docx_path)
    return result

def get_writer_titlecase_doublespace__6e2b5ade87cb088089f67dc5070eb77b_qw35sft2_362038d8(env, config: dict):
    """Download Geography_And_Magical_Realism.docx and check:
    1. Title case applied across all content paragraphs.
    2. Double line spacing on all content paragraphs.
    """
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Geography_And_Magical_Realism.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        title_case_ok = True
        for para in paras:
            text = para.text.strip()
            if not text:
                continue
            for w in text.split():
                if w and w[0].isalpha() and (not w[0].isupper()):
                    title_case_ok = False
                    break
            if not title_case_ok:
                break
        all_double = True
        content_checked = 0
        for para in paras:
            if not para.text.strip():
                continue
            content_checked += 1
            pPr = para._p.pPr
            if pPr is None:
                all_double = False
                continue
            sp = pPr.find(qn('w:spacing'))
            if sp is None:
                all_double = False
                continue
            line = sp.get(qn('w:line'))
            line_rule = sp.get(qn('w:lineRule'))
            if not (line == '480' and line_rule in ('auto', None)):
                all_double = False
        if content_checked == 0:
            all_double = False
        return {'title_case_applied': title_case_ok, 'all_double_spacing': all_double}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_subscript_title_and_body__cb7672ac12071a9e4c398d11dc4c470e_qw35sft2_96ea8482(env, config: dict):
    """
    Get formatting state of H2O_Factsheet_WA.docx:
    - Whether '2' in the Title-style paragraph has subscript formatting
    - Whether '2' in the bold-italic 'H2O—SOAK UP THE SCIENCE' run in the
      first body paragraph also has subscript formatting
    """
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/H2O_Factsheet_WA.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        subscript_in_title = False
        subscript_in_body = False
        for para in doc.paragraphs:
            if para.style.name == 'Title':
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_in_title = True
                        break
        for para in doc.paragraphs:
            if para.style.name != 'Title' and 'H2O' in para.text and ('SOAK UP THE SCIENCE' in para.text):
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_in_body = True
                        break
                break
        return {'subscript_in_title': subscript_in_title, 'subscript_in_body': subscript_in_body}
    finally:
        os.unlink(tmp_path)

def get_writer_font_italic__c3632663016de6f20637b59e83145d05_qw35sft2_4378d365(env, config: dict):
    """Read docx from VM and extract font name and italic state across all runs."""
    import tempfile
    import os
    from collections import Counter
    path = config.get('path', '/home/user/Desktop/Dublin_Zoo_Intro.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        font_names = []
        runs_with_text = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    runs_with_text.append(run)
                    if run.font.name:
                        font_names.append(run.font.name)
        most_common_font = Counter(font_names).most_common(1)[0][0] if font_names else None
        all_same_font = len(set(font_names)) <= 1 if font_names else False
        if runs_with_text:
            italic_count = sum((1 for r in runs_with_text if r.italic is True))
            all_italic = italic_count >= len(runs_with_text) * 0.9
        else:
            all_italic = False
        return {'font_name': most_common_font, 'all_same_font': all_same_font, 'all_italic': all_italic}
    finally:
        os.unlink(tmp_path)

def get_doc_page_breaks__b7c367c074a6362d7d9c85a08867a367_qw35sft2_dbc4af6a(env, config: dict):
    """Check whether the docx file contains a page break specifically before the 'Conclusion' section."""
    import tempfile, os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/The Wonders of Our Solar System.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_break': False, 'page_break_before_conclusion': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = doc.paragraphs
        has_page_break = False
        page_break_before_conclusion = False
        for para in paragraphs:
            for run in para.runs:
                breaks = run._r.findall(qn('w:br'))
                for br in breaks:
                    if br.get(qn('w:type'), 'textWrapping') == 'page':
                        has_page_break = True
            pPr = para._element.find(qn('w:pPr'))
            if pPr is not None:
                pbBefore = pPr.find(qn('w:pageBreakBefore'))
                if pbBefore is not None:
                    val = pbBefore.get(qn('w:val'), '1')
                    if val not in ('0', 'false'):
                        has_page_break = True
        for i, para in enumerate(paragraphs):
            para_text = para.text.strip()
            if para_text.lower().startswith('conclusion'):
                pPr = para._element.find(qn('w:pPr'))
                if pPr is not None:
                    pbBefore = pPr.find(qn('w:pageBreakBefore'))
                    if pbBefore is not None:
                        val = pbBefore.get(qn('w:val'), '1')
                        if val not in ('0', 'false'):
                            page_break_before_conclusion = True
                            break
                if not page_break_before_conclusion:
                    for j in range(i - 1, -1, -1):
                        candidate = paragraphs[j]
                        for run in candidate.runs:
                            breaks = run._r.findall(qn('w:br'))
                            for br in breaks:
                                br_type = br.get(qn('w:type'), 'textWrapping')
                                if br_type == 'page':
                                    page_break_before_conclusion = True
                        if page_break_before_conclusion:
                            break
                        if candidate.text.strip():
                            break
                break
        return {'has_page_break': has_page_break, 'page_break_before_conclusion': page_break_before_conclusion}
    except Exception as e:
        return {'error': str(e), 'has_page_break': False, 'page_break_before_conclusion': False}
    finally:
        os.unlink(tmp_path)

def get_docx_italic_and_title_size__792e97b1e8b831bfb11064f6cde55b00_qw35sft2_660a7e52(env, config: dict):
    """Download Y22-2119-assign4.docx and check:
    1. Whether all italic runs are 14pt.
    2. Whether the first paragraph's runs are 16pt (title heading).
    """
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Y22-2119-assign4.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        italic_sizes = []
        for para in paras:
            for run in para.runs:
                if run.italic and run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    italic_sizes.append(size_pt)
        all_italic_14 = all((s == 14.0 for s in italic_sizes)) if italic_sizes else False
        title_sizes = []
        if paras:
            for run in paras[0].runs:
                if run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    title_sizes.append(size_pt)
        all_title_16 = all((s == 16.0 for s in title_sizes)) if title_sizes else False
        return {'italic_count': len(italic_sizes), 'all_italic_size_14': all_italic_14, 'title_run_count': len(title_sizes), 'all_title_size_16': all_title_16}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_heading_align_italic__edb463affe3b938af85226bb52cb1b03_qw35sft2_b85ec2f9(env, config: dict):
    """Get heading alignment and italic status from the Constitution docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Constitution_Template_With_Guidelines.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        is_centered = para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        is_italic = any((run.italic is True for run in para.runs))
        return {'is_centered': is_centered, 'is_italic': is_italic}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_para0_italic__d2318c984412daccf01ad6c479224170_qw35sft2_1c924744(env, config: dict):
    """Get italic state of all runs in the first paragraph of the tutorial guidelines doc."""
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/CCCH9003_Tutorial_guidelines.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        if not para.runs:
            return {'italic': False, 'run_count': 0}
        all_italic = all((run.italic is True for run in para.runs))
        return {'italic': all_italic, 'run_count': len(para.runs)}
    finally:
        os.unlink(tmp_path)

def get_docx_last_para_bold__c977b8be915d8296c568265b630cf188_qw35sft2_851b1d39(env, config: dict):
    import tempfile, os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/GEOG2169_Course_Outline_2022-23.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_bold': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content_paras = [p for p in doc.paragraphs if p.text.strip()]
        if not content_paras:
            return {'error': 'No content paragraphs found', 'has_bold': False}
        last_para = content_paras[-1]
        runs_with_text = [r for r in last_para.runs if r.text.strip()]
        if not runs_with_text:
            return {'has_bold': False, 'text': last_para.text[:100]}
        all_bold = all((bool(r.bold) for r in runs_with_text))
        return {'has_bold': all_bold, 'text': last_para.text[:100]}
    finally:
        os.unlink(tmp_path)

def get_word_font_colors__3d9f327a66b9025c03e34cafe8c88fe3_qw35sft2_ef29a4bb(env, config: dict):
    """Extract font colors for all words in the table of the docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Dolch_Sight_Words_Primer.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        word_colors = {}
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        text = para.text.strip().lower()
                        if not text:
                            continue
                        color_hex = None
                        for run in para.runs:
                            try:
                                if run.font.color.rgb is not None:
                                    color_hex = str(run.font.color.rgb).upper()
                            except Exception:
                                pass
                            break
                        word_colors[text] = color_hex
        return {'word_colors': word_colors}
    except Exception as e:
        return {'error': str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_docx_spacing_italic_conclusion__bb617d210691b208d4c61100467118c1_qw35sft2_e414146c(env, config: dict):
    """Get line spacings for 3 paragraphs and italic status of the conclusion paragraph."""
    import tempfile, os
    from docx import Document
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/CCHU9045_Course_Outline_2019-20.docx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        if len(paragraphs) < 3:
            return {'error': f'Expected at least 3 non-empty paragraphs, found {len(paragraphs)}'}
        intro_spacing = paragraphs[0].paragraph_format.line_spacing
        body_spacing = paragraphs[1].paragraph_format.line_spacing
        conclusion_spacing = paragraphs[2].paragraph_format.line_spacing
        conclusion_runs = paragraphs[2].runs
        conclusion_italic = len(conclusion_runs) > 0 and all((run.italic is True for run in conclusion_runs))
        return {'intro_spacing': float(intro_spacing) if intro_spacing is not None else None, 'body_spacing': float(body_spacing) if body_spacing is not None else None, 'conclusion_spacing': float(conclusion_spacing) if conclusion_spacing is not None else None, 'conclusion_italic': conclusion_italic}
    finally:
        os.unlink(tmp_path)

def get_writer_ref_footer_state__78615d30a3700f795271d9e8e63e0686_qw35sft2_1ff19677(env, config: dict):
    """Read docx and check Steinberg presence, add-here removal, and footer page numbers."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join((p.text for p in doc.paragraphs))
        has_steinberg = 'Steinberg' in full_text
        has_add_here = '<add here>' in full_text
        has_footer_page_num = False
        for section in doc.sections:
            footer = section.footer
            footer_xml = footer._element.xml
            if 'PAGE' in footer_xml or 'NUMPAGES' in footer_xml or 'w:fldChar' in footer_xml:
                has_footer_page_num = True
                break
            for para in footer.paragraphs:
                if para.text.strip():
                    has_footer_page_num = True
                    break
            if has_footer_page_num:
                break
        return {'has_steinberg': has_steinberg, 'has_add_here_marker': has_add_here, 'has_footer_page_num': has_footer_page_num}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_page_break_count__e042d7b442613a05635c554b051c41b1_qw35sft2_7b22bea1(env, config: dict):
    """Count explicit run-level page breaks in Sample_Statutory_Declaration.docx."""
    import tempfile
    import os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/Sample_Statutory_Declaration.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        explicit_breaks = 0
        for para in doc.paragraphs:
            for run in para.runs:
                for br in run._element.findall('.//' + qn('w:br')):
                    if br.get(qn('w:type')) == 'page':
                        explicit_breaks += 1
        return {'explicit_page_breaks': explicit_breaks}
    finally:
        os.unlink(tmp_path)

def get_odt_highlight_fontsize__4ce1d37833ab4f6aa409d7454a343799_qw35sft2_cedb6e29(env, config: dict):
    _NS_TEXT = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
    _NS_STYLE = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
    _NS_FO = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
    NS_S = '{' + _NS_STYLE + '}'
    NS_T = '{' + _NS_TEXT + '}'
    NS_FO = '{' + _NS_FO + '}'
    path = config.get('path', '/home/user/Desktop/sample-recruitment-phone-script.odt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Could not load ODT file'}
    with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path) as zf:
            content = zf.read('content.xml').decode('utf-8')
            styles_xml = zf.read('styles.xml').decode('utf-8')
    finally:
        os.unlink(tmp_path)
    has_highlights = False
    try:
        root = ET.fromstring(content)
        for style in root.iter(NS_S + 'style'):
            if style.get(NS_S + 'family') != 'text':
                continue
            for tp in style.iter(NS_S + 'text-properties'):
                bg = tp.get(NS_FO + 'background-color', '')
                if bg and bg.lower() not in ('transparent', '', 'automatic', '#00000000'):
                    has_highlights = True
                    break
            if has_highlights:
                break
    except ET.ParseError:
        pass
    title_text = 'Sample Recruitment Phone Script'
    title_style_names = []
    try:
        root = ET.fromstring(content)
        for para in root.iter(NS_T + 'p'):
            all_text = ''.join(para.itertext())
            if title_text in all_text:
                para_style = para.get(NS_T + 'style-name', '')
                if para_style:
                    title_style_names.append(para_style)
                for span in para.iter(NS_T + 'span'):
                    span_style = span.get(NS_T + 'style-name', '')
                    if span_style:
                        title_style_names.append(span_style)
                break
    except ET.ParseError:
        pass
    title_font_size = None
    for sname in title_style_names:
        props = {}
        for xml_src in [content, styles_xml]:
            try:
                root = ET.fromstring(xml_src)
            except ET.ParseError:
                continue
            for style in root.iter(NS_S + 'style'):
                if style.get(NS_S + 'name') == sname:
                    for tp in style.iter(NS_S + 'text-properties'):
                        props = dict(tp.attrib)
        fs = props.get(NS_FO + 'font-size', '')
        if fs and fs.strip().endswith('pt'):
            try:
                title_font_size = float(fs.strip()[:-2])
                break
            except ValueError:
                pass
    return {'has_highlights': has_highlights, 'title_font_size': title_font_size}

def get_doc_line_spacing__80aefd2088f27d1be89724efb49ed091_qw35sft2_0db69134(env, config: dict):
    """Read line spacing values from all non-empty paragraphs in a docx file."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/04 CHIN9505 EBook Purchasing info 2021 Jan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        spacings = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            ls = para.paragraph_format.line_spacing
            if ls is not None:
                ls_val = float(ls) if not hasattr(ls, 'pt') else float(ls.pt) / 12.0
            else:
                ls_val = None
            spacings.append({'text_start': para.text[:40], 'line_spacing': ls_val})
        return {'spacings': spacings, 'count': len(spacings)}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_footer_and_header__5d6628525b72a5805637db6ac940e950_qw35sft2_bc1fbc77(env, config: dict):
    """
    Get footer page-number presence and whether a non-empty header has been added.
    """
    path = config.get('path', '/home/user/Desktop/LibreOffice_Open_Source_Word_Processing.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_field': False, 'has_header_text': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        has_page_field = False
        has_header_text = False
        with zipfile.ZipFile(tmp_path) as z:
            names = z.namelist()
            footer_files = [n for n in names if re.match('word/footer\\d*\\.xml', n)]
            for fname in footer_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                if re.search('\\bPAGE\\b', content):
                    has_page_field = True
                    break
            header_files = [n for n in names if re.match('word/header\\d*\\.xml', n)]
            for fname in header_files:
                content = z.read(fname).decode('utf-8', errors='ignore')
                texts = re.findall('<w:t[^>]*>([^<]+)</w:t>', content)
                combined = ''.join(texts).strip()
                if combined:
                    has_header_text = True
                    break
        return {'has_page_field': has_page_field, 'has_header_text': has_header_text}
    finally:
        os.unlink(tmp_path)

def get_docx_subscript_and_heading_underline__b75f28775a1b3fe6faeb633d14a05fab_qw35sft2_913900f0(env, config: dict):
    """
    Get formatting state of H2O_Factsheet_WA.docx:
    - Whether the '2' in 'H2O' in the title paragraph has subscript formatting
    - Whether the 'Fact sheet' heading paragraph has underline formatting
    """
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/H2O_Factsheet_WA.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        subscript_2_in_title = False
        heading_underlined = False
        for para in doc.paragraphs:
            if para.style.name == 'Title':
                for run in para.runs:
                    if run.font.subscript is True and '2' in run.text:
                        subscript_2_in_title = True
                        break
        for para in doc.paragraphs:
            if para.style.name == 'Heading 1' and 'Fact sheet' in para.text:
                for run in para.runs:
                    if run.underline is True:
                        heading_underlined = True
                        break
                break
        return {'subscript_2_in_title': subscript_2_in_title, 'heading_underlined': heading_underlined}
    finally:
        os.unlink(tmp_path)

def get_writer_titlecase_bold__2777e85b511814778d7406b21395647f_qw35sft2_643ea196(env, config: dict):
    """Download Geography_And_Magical_Realism.docx and check:
    1. Title case applied across all content paragraphs.
    2. Bold formatting on the title paragraph (para 0).
    """
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Geography_And_Magical_Realism.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paras = doc.paragraphs
        title_case_ok = True
        for para in paras:
            text = para.text.strip()
            if not text:
                continue
            for w in text.split():
                if w and w[0].isalpha() and (not w[0].isupper()):
                    title_case_ok = False
                    break
            if not title_case_ok:
                break
        title_para = paras[0]
        title_bold = False
        if title_para.runs:
            title_bold = all((run.bold is True for run in title_para.runs))
        return {'title_case_applied': title_case_ok, 'title_bold': title_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_footer_page_numbers__699ab5651848f548e06466de9777d875_qw35sft2_dbe640d9(env, config: dict):
    """Check whether the docx file footer contains page number fields."""
    import tempfile, os
    from docx import Document
    from docx.oxml.ns import qn
    path = config.get('path', '/home/user/Desktop/The Wonders of Our Solar System.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'has_page_numbers': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        has_page_numbers = False
        for section in doc.sections:
            footer = section.footer
            if footer is None:
                continue
            footer_xml = footer._element.xml
            if 'PAGE' in footer_xml or 'pgNum' in footer_xml or 'w:fldChar' in footer_xml:
                for para in footer.paragraphs:
                    for run in para.runs:
                        instr_texts = run._r.findall(qn('w:instrText'))
                        for instr in instr_texts:
                            if instr.text and 'PAGE' in instr.text.upper():
                                has_page_numbers = True
                        fld_chars = run._r.findall(qn('w:fldChar'))
                        if fld_chars:
                            has_page_numbers = True
                if footer._element.findall('.//' + qn('w:pgNum')):
                    has_page_numbers = True
                for para in footer.paragraphs:
                    for run in para.runs:
                        if run._r.findall('.//' + qn('w:pgNum')):
                            has_page_numbers = True
        return {'has_page_numbers': has_page_numbers}
    except Exception as e:
        return {'error': str(e), 'has_page_numbers': False}
    finally:
        os.unlink(tmp_path)

def get_docx_dedup_state__f922eeeed3d49013fd1e13103dbfb120_qw35sft2_7be396a3(env, config: dict):
    """Check deduplication state of a docx: unique train IDs and whether each appears only once."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not available'}
    path = config.get('path', '/home/user/Desktop/HK_train_record.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        from collections import Counter
        train_counts = Counter()
        valid_lines = 0
        for l in lines:
            parts = l.split(',')
            if len(parts) >= 2:
                train_id = parts[1].strip()
                train_counts[train_id] += 1
                valid_lines += 1
        unique_trains = len(train_counts)
        all_unique = all((v == 1 for v in train_counts.values()))
        return {'total_lines': len(lines), 'valid_lines': valid_lines, 'unique_train_count': unique_trains, 'all_unique': all_unique, 'train_counts': dict(train_counts)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_and_footer__d1ae06a2cb3f08c6662cf614ce58e285_qw35sft2_25efcd2d(env, config: dict):
    """
    Fetch LibreOffice XCU config and the docx file, then return:
    - default_font: the Standard font name from XCU
    - has_page_numbers: whether page numbers appear in the document footer
    """
    result = {'default_font': None, 'has_page_numbers': False}
    xcu_bytes = env.controller.get_file('/home/user/.config/libreoffice/4/user/registrymodifications.xcu')
    if xcu_bytes:
        with tempfile.NamedTemporaryFile(suffix='.xcu', delete=False) as f:
            f.write(xcu_bytes)
            xcu_path = f.name
        try:
            tree = ET.parse(xcu_path)
            root = tree.getroot()
            ns = {'oor': 'http://openoffice.org/2001/registry'}
            for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', ns):
                for prop in elem.findall('.//prop[@oor:name="Standard"]', ns):
                    for value in prop.findall('value', ns):
                        result['default_font'] = value.text
        except Exception:
            pass
        finally:
            os.unlink(xcu_path)
    docx_bytes = env.controller.get_file('/home/user/Desktop/loa-one-time-submission-sealand.docx')
    if docx_bytes:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as f:
            f.write(docx_bytes)
            docx_path = f.name
        try:
            from docx import Document
            doc = Document(docx_path)
            has_pn = False
            for section in doc.sections:
                footer = section.footer
                if footer is None:
                    continue
                footer_text = footer.paragraphs[0].text if footer.paragraphs else ''
                if any((char.isdigit() for char in footer_text)):
                    has_pn = True
                    break
                from docx.oxml.ns import qn
                footer_xml = footer._element.xml if hasattr(footer, '_element') else ''
                if 'PAGE' in footer_xml or 'fldChar' in footer_xml:
                    has_pn = True
                    break
            result['has_page_numbers'] = has_pn
        except Exception:
            pass
        finally:
            os.unlink(docx_path)
    return result

def get_docx_italic_bold_size__57bf430ce2e461a40bf234942a8ba4ad_qw35sft2_fe2c4902(env, config: dict):
    """Download Y22-2119-assign4.docx and check whether all italic runs are 14pt and bold."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Y22-2119-assign4.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        italic_runs = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.italic and run.text.strip():
                    size = run.font.size
                    size_pt = size.pt if size else None
                    italic_runs.append({'size': size_pt, 'bold': bool(run.bold)})
        if not italic_runs:
            return {'italic_count': 0, 'all_italic_size_14': False, 'all_italic_bold': False}
        all_14 = all((r['size'] == 14.0 for r in italic_runs))
        all_bold = all((r['bold'] for r in italic_runs))
        return {'italic_count': len(italic_runs), 'all_italic_size_14': all_14, 'all_italic_bold': all_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_writer_font_size__5934b0de41c172c0d1662adc3bbc874f_qw35sft2_a3711a32(env, config: dict):
    """Read docx from VM and extract font name and font size stats across all runs."""
    import tempfile
    import os
    from collections import Counter
    path = config.get('path', '/home/user/Desktop/Dublin_Zoo_Intro.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        font_names = []
        font_sizes = []
        for para in doc.paragraphs:
            for run in para.runs:
                if run.text.strip():
                    if run.font.name:
                        font_names.append(run.font.name)
                    if run.font.size:
                        font_sizes.append(round(run.font.size.pt, 1))
        most_common_font = Counter(font_names).most_common(1)[0][0] if font_names else None
        most_common_size = Counter(font_sizes).most_common(1)[0][0] if font_sizes else None
        all_same_font = len(set(font_names)) <= 1 if font_names else False
        all_same_size = len(set(font_sizes)) <= 1 if font_sizes else False
        return {'font_name': most_common_font, 'font_size_pt': most_common_size, 'all_same_font': all_same_font, 'all_same_size': all_same_size}
    finally:
        os.unlink(tmp_path)

def get_writer_heading_align_body_bold__e1c39faa270c25698b597d87598bba49_qw35sft2_7ed4bddf(env, config: dict):
    """Get heading alignment and body paragraph bold status from the Constitution docx."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_path = config.get('path', '/home/user/Desktop/Constitution_Template_With_Guidelines.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if len(doc.paragraphs) < 3:
            return {'error': 'Not enough paragraphs'}
        heading_para = doc.paragraphs[0]
        body_para = doc.paragraphs[2]
        heading_centered = heading_para.alignment == WD_ALIGN_PARAGRAPH.CENTER
        non_empty_runs = [run for run in body_para.runs if run.text.strip()]
        if non_empty_runs:
            body_bold = all((run.bold is True for run in non_empty_runs))
        else:
            body_bold = False
        return {'heading_centered': heading_centered, 'body_bold': body_bold}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_para0_font_size__7211ff5a62d6ff910c67632b31846d1c_qw35sft2_982eaf9d(env, config: dict):
    """Get font size (in points) of the first run in the first paragraph of the tutorial guidelines doc."""
    import tempfile
    import os
    from docx import Document
    file_bytes = env.controller.get_file('/home/user/Desktop/CCCH9003_Tutorial_guidelines.docx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.paragraphs:
            return {'error': 'No paragraphs found'}
        para = doc.paragraphs[0]
        if not para.runs:
            return {'font_size_pt': None}
        size_emu = para.runs[0].font.size
        if size_emu is None:
            return {'font_size_pt': None}
        size_pt = size_emu / 12700.0
        return {'font_size_pt': size_pt}
    finally:
        os.unlink(tmp_path)

def get_docx_second_para_italic__aa396abb2d9562e18e97067a7c8c6fe9_qw35sft2_ddda8fe7(env, config: dict):
    import tempfile, os
    from docx import Document
    file_path = config.get('path', '/home/user/Desktop/GEOG2169_Course_Outline_2022-23.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'has_italic': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content_paras = [p for p in doc.paragraphs if p.text.strip()]
        if len(content_paras) < 2:
            return {'error': 'Not enough paragraphs', 'has_italic': False}
        second_para = content_paras[1]
        runs_with_text = [r for r in second_para.runs if r.text.strip()]
        if not runs_with_text:
            return {'has_italic': False, 'text': second_para.text[:100]}
        all_italic = all((bool(r.italic) for r in runs_with_text))
        return {'has_italic': all_italic, 'text': second_para.text[:100]}
    finally:
        os.unlink(tmp_path)

def get_docx_spacing_center_intro__5f00d786b29854e31074fa3cfeefea1b_qw35sft2_3e7fe962(env, config: dict):
    """Get line spacings for 3 paragraphs and alignment of the introduction."""
    import tempfile, os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/CCHU9045_Course_Outline_2019-20.docx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        if len(paragraphs) < 3:
            return {'error': f'Expected at least 3 non-empty paragraphs, found {len(paragraphs)}'}
        intro_spacing = paragraphs[0].paragraph_format.line_spacing
        body_spacing = paragraphs[1].paragraph_format.line_spacing
        conclusion_spacing = paragraphs[2].paragraph_format.line_spacing
        intro_alignment = paragraphs[0].alignment
        intro_centered = intro_alignment == WD_ALIGN_PARAGRAPH.CENTER
        return {'intro_spacing': float(intro_spacing) if intro_spacing is not None else None, 'body_spacing': float(body_spacing) if body_spacing is not None else None, 'conclusion_spacing': float(conclusion_spacing) if conclusion_spacing is not None else None, 'intro_centered': intro_centered}
    finally:
        os.unlink(tmp_path)

def get_writer_citation14_state__21b44c209dfc410a39a6375bc0042826_qw35sft2_983282d0(env, config: dict):
    """Read docx and check for Steinberg reference and [14] citation in body paragraphs."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed'}
    path = config.get('path', '/home/user/Desktop/Essay_Writing_English_for_uni.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = doc.paragraphs
        body_text = ''
        refs_text = ''
        in_refs = False
        for p in paragraphs:
            if p.text.strip() == 'References' or p.text.strip() == 'Reference List':
                in_refs = True
            if in_refs:
                refs_text += p.text + '\n'
            else:
                body_text += p.text + '\n'
        has_steinberg = 'Steinberg' in refs_text or 'Steinberg' in body_text
        has_bracket14_in_body = '[14]' in body_text
        has_add_here = '<add here>' in body_text
        return {'has_steinberg': has_steinberg, 'has_bracket14_in_body': has_bracket14_in_body, 'has_add_here_marker': has_add_here}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_invoiceGES_in_problematic__cea86e2d544f9d987b2adcf034f36c4a_qw35sft2_7a0e7baa(env, config: dict):
    """Check if 'Invoice # GES-20220215-82.pdf' is inside Desktop/problematic folder."""
    script = 'test -d "/home/user/Desktop/problematic" && echo "folder_ok" || echo "no_folder"; test -f "/home/user/Desktop/problematic/Invoice # GES-20220215-82.pdf" && echo "file_ok" || echo "no_file"'
    result = env.controller.run_bash_script(script, timeout=10)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    lines = output.splitlines()
    return {'folder_exists': 'folder_ok' in lines, 'file_in_problematic': 'file_ok' in lines}

def get_docx_with_header__2354b786c1a1e08afc94d3d722cbb7a6_qw35sft2_bdd3f828(env, config: dict):
    """Read notes.docx and check for header line + notes content."""
    path = config.get('path', '/home/user/Desktop/notes.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'lines': [], 'first_line': None}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from docx import Document
        doc = Document(tmp_path)
        lines = [para.text for para in doc.paragraphs if para.text.strip()]
        first_line = lines[0] if lines else None
        return {'lines': lines, 'first_line': first_line, 'count': len(lines)}
    except Exception as e:
        return {'error': str(e), 'lines': [], 'first_line': None}
    finally:
        os.unlink(tmp_path)

def get_main_py_first_line__7b62591361b18caf2ed6aa916d0c9cf5_qw35sft2_61a5d212(env, config: dict):
    """Read main.py from the project directory and return its first line."""
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/project/main.py')
        raw = file_bytes.decode('utf-8') if isinstance(file_bytes, bytes) else str(file_bytes)
        lines = raw.splitlines()
        return lines[0].strip() if lines else ''
    except Exception:
        return ''

def get_docx_text__ed93db9128130b7a782216ac4507a0ca_qw35sft2_c07f733b(env, config: dict):
    """Read full text content from a .docx file via the VM controller."""
    import tempfile
    import os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed', 'full_text': ''}
    path = config.get('path', '/home/user/Desktop/Answer.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'full_text': ''}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join((p.text for p in doc.paragraphs))
        return {'full_text': full_text}
    except Exception as e:
        return {'error': str(e), 'full_text': ''}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_docx_first_para__40a84e6a55059151959a2b459d4c099e_qw35sft2_89c54b9d(env, config: dict):
    """Read the first non-empty paragraph from a docx file on the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        import subprocess
        subprocess.run(['pip', 'install', 'python-docx'], capture_output=True)
        from docx import Document
    path = config.get('path', '/home/user/Desktop/students work/case study.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for para in doc.paragraphs:
            if para.text.strip():
                return {'first_paragraph': para.text.strip()}
        return {'first_paragraph': ''}
    finally:
        os.unlink(tmp_path)

def get_sar_disk_report_state__fc8ac0599dbb02ca56427b3ae265c796_qw35sft2_6bd8bd41(env, config: dict):
    """Read Disk_IO_Report.txt from Desktop and return file info."""
    result = env.controller.run_bash_script('python3 -c "import os; path = os.path.expanduser(\'~/Desktop/Disk_IO_Report.txt\'); exists = os.path.isfile(path); has_disk_data = False; line_count = 0; if exists:     with open(path) as f: content = f.read();     lines = content.splitlines();     line_count = len(lines);     has_disk_data = any(\'tps\' in l.lower() or \'rkb/s\' in l.lower() or \'wkb/s\' in l.lower() or \'DEV\' in l for l in lines); print(f\'exists={exists},has_disk_data={has_disk_data},line_count={line_count}\'); "', timeout=30)
    if not result or result.get('returncode', 1) != 0:
        return {'error': 'script failed', 'exists': False, 'has_disk_data': False, 'line_count': 0}
    output = result.get('output', '').strip()
    data = {}
    for part in output.split(','):
        k, _, v = part.partition('=')
        data[k.strip()] = v.strip()
    return {'exists': data.get('exists', 'False') == 'True', 'has_disk_data': data.get('has_disk_data', 'False') == 'True', 'line_count': int(data.get('line_count', 0))}

def get_docx_writer_state__8ae851f2fbec3e62c3a4b8c28bfb6c8b_qw35sft2_a3812e2a(env, config: dict):
    """Read docx from VM and extract tables, paragraphs, and ordered body elements."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed', 'tables': [], 'paragraphs': [], 'table_count': 0, 'elements': []}
    path = config.get('path', '/home/user/Documents/awesome-desktop/awe_desk_env.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'tables': [], 'paragraphs': [], 'table_count': 0, 'elements': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            tables.append(rows)
        paragraphs = [para.text.strip() for para in doc.paragraphs]
        elements = []
        table_idx = 0
        para_idx = 0
        for child in doc.element.body:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'tbl':
                if table_idx < len(tables):
                    elements.append({'type': 'table', 'index': table_idx, 'data': tables[table_idx]})
                    table_idx += 1
            elif tag == 'p':
                if para_idx < len(paragraphs):
                    elements.append({'type': 'paragraph', 'text': paragraphs[para_idx]})
                    para_idx += 1
        return {'table_count': len(tables), 'tables': tables, 'paragraphs': paragraphs, 'has_table': len(tables) > 0, 'elements': elements}
    except Exception as e:
        return {'error': str(e), 'tables': [], 'paragraphs': [], 'table_count': 0, 'elements': []}
    finally:
        os.unlink(tmp_path)

def get_settings_snake_size__85e39531da2c848a1afaf10e25ba3dad_qw35sft2_f06e060f(env, config: dict):
    """Read settings.py and extract the SNAKE_SIZE value."""
    import re
    file_bytes = env.controller.get_file('/home/user/Desktop/snake/settings.py')
    if not file_bytes:
        return {'error': 'settings.py not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    match = re.search('^SNAKE_SIZE\\s*=\\s*(\\d+)', content, re.MULTILINE)
    if match:
        return {'snake_size': int(match.group(1))}
    return {'error': 'SNAKE_SIZE not found in settings.py'}

def get_invoiceTII_in_problematic__7ecb63a4609903641df39f4754b7248f_qw35sft2_84d08576(env, config: dict):
    """Check if 'invoice TII-20220301-90.pdf' is inside Desktop/problematic folder."""
    script = 'test -d "/home/user/Desktop/problematic" && echo "folder_ok" || echo "no_folder"; test -f "/home/user/Desktop/problematic/invoice TII-20220301-90.pdf" && echo "file_ok" || echo "no_file"'
    result = env.controller.run_bash_script(script, timeout=10)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    lines = output.splitlines()
    return {'folder_exists': 'folder_ok' in lines, 'file_in_problematic': 'file_ok' in lines}

def get_book_copy_in_documents__55d3e622f74c4c6d9051e6e358a64ecd_qw35sft2_a3d090d3(env, config: dict):
    """Check if the book PDF was copied to the Documents folder."""
    dest_path = config.get('dest_path', '/home/user/Documents/Spectral Graph Theory.pdf')
    result = env.controller.run_bash_script(f'test -f "{dest_path}" && echo "exists" || echo "not_found"', timeout=15)
    output = result.get('output', '') if isinstance(result, dict) else str(result)
    return {'file_exists': 'exists' in output}

def get_combined_state__3001e2e091ccfaef8fccec23c5f15501_qw35sft2_9a14dca8(env, config: dict):
    """Get VS Code extension list and resized.png dimensions.

    Returns a dict with:
        extensions (str): newline-separated list of installed VS Code extensions
        img_width  (int|None): width of /home/user/Desktop/resized.png, or None
        img_height (int|None): height of /home/user/Desktop/resized.png, or None
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    try:
        ext_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        ext_output = ext_resp.json().get('output', '') if ext_resp.status_code == 200 else ''
    except Exception:
        ext_output = ''
    img_width, img_height = (None, None)
    tmp_path = None
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/resized.png')
        if file_bytes:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            from PIL import Image
            img = Image.open(tmp_path)
            img_width, img_height = img.size
    except Exception:
        pass
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except Exception:
                pass
    return {'extensions': ext_output, 'img_width': img_width, 'img_height': img_height}

def get_tally_book_amount_sum__289ae9cd25fb72b631d3d5c97c4b9f82_qw35sft2_fb6205ee(env, config: dict):
    """Read tally_book.xlsx and check for PDF receipt existence in receipts folder."""
    import tempfile, os, openpyxl
    pdf_exists = False
    try:
        pdf_bytes = env.controller.get_file('/home/user/Documents/Finance/receipts/aws-invoice-2312.pdf')
        pdf_exists = pdf_bytes is not None and len(pdf_bytes) > 0
    except Exception:
        pdf_exists = False
    path = config.get('path', '/home/user/Documents/Finance/tally_book.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'total': None, 'row_count': 0, 'pdf_exists': pdf_exists}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        total = 0.0
        row_count = 0
        for row in ws.iter_rows(min_row=2, values_only=True):
            val = row[2]
            if val is not None:
                try:
                    total += float(val)
                    row_count += 1
                except (TypeError, ValueError):
                    pass
        return {'total': round(total, 4), 'row_count': row_count, 'pdf_exists': pdf_exists}
    finally:
        os.unlink(tmp_path)

def get_paper03_and_year__d508d903f596b0ecb03360bdda6892d4_qw35sft2_17558f40(env, config: dict):
    """Check /home/user/paper03.pdf exists and read /home/user/paper03_year.txt."""
    pdf_result = env.controller.run_bash_script('test -f /home/user/paper03.pdf && echo "exists" || echo "missing"', timeout=15)
    pdf_exists = False
    if pdf_result and pdf_result.get('stdout', '').strip() == 'exists':
        pdf_exists = True
    year_result = env.controller.run_bash_script('cat /home/user/paper03_year.txt 2>/dev/null', timeout=15)
    year_content = ''
    if year_result:
        year_content = year_result.get('stdout', '').strip()
    return {'pdf_exists': pdf_exists, 'year_content': year_content}

def get_docx_duration_line__48db15ca4dac03d9c5bae1d76e883852_qw35sft2_7725eb8d(env, config: dict):
    """Read the 'Duration' line from a docx file on the VM."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        import subprocess
        subprocess.run(['pip', 'install', 'python-docx'], capture_output=True)
        from docx import Document
    path = config.get('path', '/home/user/Desktop/Public Lecture Teaching Plan.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        for para in doc.paragraphs:
            text = para.text.strip()
            if text.lower().startswith('duration'):
                return {'duration_line': text}
        return {'duration_line': ''}
    finally:
        os.unlink(tmp_path)

def get_desktop_listing__00bb03d81e57baa40bd1c86ce0b1574d_qw35sft2_5b6e050b(env, config: dict):
    """List files on the VM Desktop to verify file existence."""
    result = env.controller.run_bash_script('ls /home/user/Desktop/', timeout=30)
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    return {'output': output}

def get_sar_cpu_report_state__5dfc721833abb9e182bb18218c19d630_qw35sft2_32121cc6(env, config: dict):
    """Read System_Resources_Report.txt from Desktop and return file info."""
    result = env.controller.run_bash_script('python3 -c "import os; path = os.path.expanduser(\'~/Desktop/System_Resources_Report.txt\'); exists = os.path.isfile(path); lines = []; cpu_lines = 0; has_header = False; if exists:     with open(path) as f: content = f.read();     lines = content.splitlines();     cpu_lines = sum(1 for l in lines if \' all \' in l);     has_header = any(\'%user\' in l for l in lines); print(f\'exists={exists},cpu_lines={cpu_lines},has_header={has_header}\'); "', timeout=30)
    if not result or result.get('returncode', 1) != 0:
        return {'error': 'script failed', 'exists': False, 'cpu_lines': 0, 'has_header': False}
    output = result.get('output', '').strip()
    data = {}
    for part in output.split(','):
        k, _, v = part.partition('=')
        data[k.strip()] = v.strip()
    return {'exists': data.get('exists', 'False') == 'True', 'cpu_lines': int(data.get('cpu_lines', 0)), 'has_header': data.get('has_header', 'False') == 'True'}

def get_docx_gemini_paragraphs__abfc7551de3c0f45f173130d54033329_qw35sft2_a085b892(env, config: dict):
    """Get paragraph texts from gemini_results.docx on the Desktop."""
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/gemini_results.docx')
        if not file_bytes:
            return {'error': 'file not found'}
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from docx import Document
            doc = Document(tmp_path)
            non_empty = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            full_text = '\n'.join(non_empty)
            return {'paragraph_count': len(non_empty), 'full_text': full_text, 'paragraphs': non_empty}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_writer_state__5d4910d5e528b018db8af7469dad985a_qw35sft2_4b976908(env, config: dict):
    """Read docx from VM and extract tables and paragraphs."""
    import tempfile, os
    try:
        from docx import Document
    except ImportError:
        return {'error': 'python-docx not installed', 'tables': [], 'paragraphs': [], 'table_count': 0}
    path = config.get('path', '/home/user/Documents/awesome-desktop/awe_desk_env.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'tables': [], 'paragraphs': [], 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            tables.append(rows)
        paragraphs = [para.text.strip() for para in doc.paragraphs]
        return {'table_count': len(tables), 'tables': tables, 'paragraphs': paragraphs, 'has_table': len(tables) > 0}
    except Exception as e:
        return {'error': str(e), 'tables': [], 'paragraphs': [], 'table_count': 0}
    finally:
        os.unlink(tmp_path)

def get_settings_fps__751d0a9ef75d8ec2c4c27e57248ad37c_qw35sft2_74691e56(env, config: dict):
    """Read settings.py and extract the FPS value."""
    import re
    file_bytes = env.controller.get_file('/home/user/Desktop/snake/settings.py')
    if not file_bytes:
        return {'error': 'settings.py not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    match = re.search('^FPS\\s*=\\s*(\\d+)', content, re.MULTILINE)
    if match:
        return {'fps': int(match.group(1))}
    return {'error': 'FPS not found in settings.py'}

def get_pandoc_install_status__79ba7b4d76657204de8c0df9dc21c02e_qw35sft2_452b3542(env, config: dict):
    """Check whether pandoc is installed by running 'which pandoc'."""
    try:
        result_raw = env.controller.run_bash_script('which pandoc 2>/dev/null', timeout=30)
        if isinstance(result_raw, dict):
            output = result_raw.get('output', result_raw.get('stdout', '')).strip()
        elif isinstance(result_raw, str):
            output = result_raw.strip()
        else:
            output = ''
        installed = bool(output) and 'pandoc' in output
        return {'installed': installed, 'path': output}
    except Exception as e:
        return {'installed': False, 'path': '', 'error': str(e)}

def get_invoice243729_in_problematic__da5edbd81a2d596fb7eb6ef7b52c2151_qw35sft2_e46fc3b3(env, config: dict):
    """Check if Invoice # 243729.pdf is inside Desktop/problematic folder."""
    script = 'test -d "/home/user/Desktop/problematic" && echo "folder_ok" || echo "no_folder"; test -f "/home/user/Desktop/problematic/Invoice # 243729.pdf" && echo "file_ok" || echo "no_file"'
    result = env.controller.run_bash_script(script, timeout=10)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    lines = output.splitlines()
    return {'folder_exists': 'folder_ok' in lines, 'file_in_problematic': 'file_ok' in lines}

def get_conda_install_state__d5abdc9092a51546a82222c972d5edf6_qw35sft2_98581235(env, config: dict):
    """Get conda binary presence and .bashrc conda init configuration state."""
    result = {}
    conda_check = env.controller.run_bash_script('test -f /home/user/miniconda3/bin/conda && echo conda_found || echo conda_not_found', timeout=15)
    conda_out = (conda_check.get('output', '') or '').strip()
    result['conda_installed'] = 'conda_found' in conda_out
    bashrc_check = env.controller.run_bash_script('grep -c "conda initialize" /home/user/.bashrc 2>/dev/null || echo 0', timeout=10)
    bashrc_out = (bashrc_check.get('output', '') or '').strip()
    try:
        count = int(bashrc_out)
    except (ValueError, TypeError):
        count = 0
    result['bashrc_configured'] = count > 0
    return result

def get_paper01_and_count__101069bcaae5b37861b743495ae4859d_qw35sft2_6d9a0ea7(env, config: dict):
    """Check /home/user/paper01.pdf exists and read /home/user/paper_count.txt."""
    pdf_result = env.controller.run_bash_script('test -f /home/user/paper01.pdf && echo "exists" || echo "missing"', timeout=15)
    pdf_exists = False
    if pdf_result and pdf_result.get('stdout', '').strip() == 'exists':
        pdf_exists = True
    count_result = env.controller.run_bash_script('cat /home/user/paper_count.txt 2>/dev/null', timeout=15)
    count_content = ''
    if count_result:
        count_content = count_result.get('stdout', '').strip()
    return {'pdf_exists': pdf_exists, 'count_content': count_content}

def get_desktop_listing__b0e8ccdfade877a9b91932809683825c_qw35sft2_bfc0cbaa(env, config: dict):
    """List files on the VM Desktop to verify file existence."""
    result = env.controller.run_bash_script('ls /home/user/Desktop/', timeout=30)
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    return {'output': output}

def get_tally_book_last_row__97f6f4530ffb08062fee8d0568f59984_qw35sft2_3aa47332(env, config: dict):
    """Read tally_book.xlsx last row and check PDF file existence in receipts folder."""
    import tempfile, os, openpyxl
    pdf_path = '/home/user/Documents/Finance/receipts/aws-invoice-2312.pdf'
    pdf_bytes = env.controller.get_file(pdf_path)
    pdf_found = bool(pdf_bytes)
    path = config.get('path', '/home/user/Documents/Finance/tally_book.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'pdf_found': pdf_found}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        max_row = ws.max_row
        if max_row < 2:
            return {'error': 'No data rows found', 'pdf_found': pdf_found}
        service = ws.cell(row=max_row, column=1).value
        month = ws.cell(row=max_row, column=2).value
        amount = ws.cell(row=max_row, column=3).value
        return {'pdf_found': pdf_found, 'service': str(service).strip() if service is not None else None, 'month': float(month) if month is not None else None, 'amount': float(amount) if amount is not None else None, 'row': max_row}
    finally:
        os.unlink(tmp_path)

def get_copy_move_state__ec373221258728f3163857f2bdfd353a_qw35sft2_d2ffe179(env, config: dict):
    """Check file1 in dir1/dir2/dir3 and absence of file1 in home dir."""
    result = env.controller.run_bash_script('echo dir1:$(test -f /home/user/dir1/file1 && echo yes || echo no)\necho dir2:$(test -f /home/user/dir2/file1 && echo yes || echo no)\necho dir3:$(test -f /home/user/dir3/file1 && echo yes || echo no)\necho orig:$(test -f /home/user/file1 && echo yes || echo no)', timeout=10)
    if not result:
        return {'error': 'command failed'}
    output = result.get('output', '')
    return {'dir1_has_file1': 'dir1:yes' in output, 'dir2_has_file1': 'dir2:yes' in output, 'dir3_has_file1': 'dir3:yes' in output, 'original_removed': 'orig:no' in output}

def get_clock_dual_settings__443802668a0bfacbf0757782ecb2ad43_qw35sft2_39fa83f8(env, config: dict):
    """Get both clock-format and clock-show-weekday GNOME settings."""
    try:
        res_format = env.controller.run_bash_script('gsettings get org.gnome.desktop.interface clock-format', timeout=10)
        res_weekday = env.controller.run_bash_script('gsettings get org.gnome.desktop.interface clock-show-weekday', timeout=10)
        clock_format = (res_format.get('output', '') or res_format.get('stdout', '')).strip() if isinstance(res_format, dict) else str(res_format).strip()
        clock_show_weekday = (res_weekday.get('output', '') or res_weekday.get('stdout', '')).strip() if isinstance(res_weekday, dict) else str(res_weekday).strip()
        return {'clock_format': clock_format, 'clock_show_weekday': clock_show_weekday}
    except Exception as e:
        return {'error': str(e)}

def get_accessibility_large_text_no_animations__855f455c8e49beda6b556868bdd9b753_qw35sft2_c390975a(env, config: dict):
    """Get large-text toggle and enable-animations states via gsettings."""
    vm_ip = env.vm_ip
    port = env.server_port
    result = {}
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.interface', 'large-text'], 'shell': False})
    if resp1.status_code == 200:
        result['large_text'] = resp1.json()['output'].strip() == 'true'
    else:
        result['large_text'] = None
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.interface', 'enable-animations'], 'shell': False})
    if resp2.status_code == 200:
        result['animations_enabled'] = resp2.json()['output'].strip() == 'true'
    else:
        result['animations_enabled'] = None
    return result

def get_gnome_favorites__41851786639af4a56b165470f26fecab_qw35sft2_bbbb11b6(env, config: dict):
    """Get the current GNOME shell favorite-apps list as a parsed Python list."""
    vm_ip = env.vm_ip
    port = env.server_port
    command = ['gsettings', 'get', 'org.gnome.shell', 'favorite-apps']
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
    if response.status_code == 200:
        apps_str = response.json()['output'].strip()
        try:
            apps = eval(apps_str)
            return {'apps': apps}
        except Exception:
            return {'error': f'Cannot parse favorites string: {apps_str}'}
    return {'error': 'Command execution failed'}

def get_volume_level__d6b3f95410b1cde5a46558c01aaafc46_qw35sft2_5eb74f22(env, config: dict):
    """Get current PulseAudio sink volume output string."""
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pactl', 'get-sink-volume', '@DEFAULT_SINK@'], 'shell': False})
    if response.status_code == 200:
        return response.json().get('output', '')
    return ''

def get_timezone_and_clock__c55229b53381e6587f722d00658cbd02_qw35sft2_79b14818(env, config: dict):
    """Get system timezone offset and GNOME clock format setting.

    Returns a dict with:
      - 'timezone_offset': string like '+0000' from `date +%z`
      - 'clock_format': string like "'24h'" or "'12h'" from gsettings
    """
    vm_ip = env.vm_ip
    port = env.server_port
    offset = ''
    clock_fmt = ''
    try:
        r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['date', '+%z'], 'shell': False}, timeout=10)
        if r1.status_code == 200:
            offset = r1.json().get('output', '').strip()
    except Exception as e:
        logger_qw35sft2_a4e44f.warning(f'Failed to get timezone offset: {e}')
    try:
        r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.interface', 'clock-format'], 'shell': False}, timeout=10)
        if r2.status_code == 200:
            clock_fmt = r2.json().get('output', '').strip()
    except Exception as e:
        logger_qw35sft2_a4e44f.warning(f'Failed to get clock format: {e}')
    return {'timezone_offset': offset, 'clock_format': clock_fmt}

def get_rename_and_sibling__9cc29f7146497c596e333ae0ebaf5848_qw35sft2_82e8f5b4(env, config: dict):
    """Check if Desktop has todo_list_Jan_2 (renamed) and todo_list_Jan_3 (newly created)."""
    vm_ip = env.vm_ip
    port = env.server_port
    r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_2 && echo yes || echo no', 'shell': True})
    renamed = r1.status_code == 200 and r1.json().get('output', '').strip() == 'yes'
    r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_3 && echo yes || echo no', 'shell': True})
    created = r2.status_code == 200 and r2.json().get('output', '').strip() == 'yes'
    r3 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_1 && echo yes || echo no', 'shell': True})
    old_absent = r3.status_code == 200 and r3.json().get('output', '').strip() == 'no'
    return {'renamed': renamed, 'created': created, 'old_absent': old_absent}

def get_notif_clock__accf3abf7a1d21f7f68eefd98414a494_qw35sft2_2b863f6d(env, config: dict):
    """Get DND (show-banners) and clock-format settings."""
    r1 = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-banners', timeout=15)
    show_banners = (r1.get('output', r1.get('stdout', '')) if isinstance(r1, dict) else str(r1)).strip()
    r2 = env.controller.run_bash_script('gsettings get org.gnome.desktop.interface clock-format', timeout=15)
    clock_format = (r2.get('output', r2.get('stdout', '')) if isinstance(r2, dict) else str(r2)).strip()
    return {'show_banners': show_banners, 'clock_format': clock_format}

def get_power_dim_and_blank__4a865a2b73390cd1be0bc5f929d03ff5_qw35sft2_fd6f3625(env, config: dict):
    """Get idle-dim and idle-delay settings from GNOME power configuration."""
    r_dim = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power idle-dim', timeout=10)
    if isinstance(r_dim, dict):
        idle_dim_raw = r_dim.get('output', r_dim.get('stdout', '')).strip()
    else:
        idle_dim_raw = str(r_dim).strip()
    r_delay = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    if isinstance(r_delay, dict):
        idle_delay_raw = r_delay.get('output', r_delay.get('stdout', '')).strip()
    else:
        idle_delay_raw = str(r_delay).strip()
    idle_dim = idle_dim_raw.lower().strip()
    parts = idle_delay_raw.split()
    idle_delay = int(parts[-1]) if parts else -1
    return {'idle_dim': idle_dim, 'idle_delay': idle_delay}

def get_php_stats__a63cdece1fd90a60a3e737af233f1764_qw35sft2_2b1cee42(env, config: dict):
    """Read php_file_count.txt and php_line_count.txt from the VM home directory."""
    out_file = env.controller.run_bash_script('cat /home/user/php_file_count.txt 2>/dev/null', timeout=30)
    if isinstance(out_file, dict):
        file_count = (out_file.get('stdout') or '').strip()
    else:
        file_count = str(out_file).strip() if out_file else ''
    out_line = env.controller.run_bash_script('cat /home/user/php_line_count.txt 2>/dev/null', timeout=30)
    if isinstance(out_line, dict):
        line_count = (out_line.get('stdout') or '').strip()
    else:
        line_count = str(out_line).strip() if out_line else ''
    return {'file_count': file_count, 'line_count': line_count}

def get_move_failed_notebooks__96695aa3604bded58ed2ec7f3c72c8de_qw35sft2_932f1c29(env, config: dict):
    """Check ./fails has moved *failed.ipynb and originals are deleted."""
    vm_ip = env.vm_ip
    port = env.server_port
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': "find /home/user/test_environment/fails -name '*failed.ipynb' -type f 2>/dev/null | sort", 'shell': True})
    fails_files = resp1.json().get('output', '').strip() if resp1.status_code == 200 else ''
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': "find /home/user/test_environment -not -path '*/fails/*' -name '*failed.ipynb' -type f 2>/dev/null | sort", 'shell': True})
    orig_files = resp2.json().get('output', '').strip() if resp2.status_code == 200 else ''
    return {'fails_files': fails_files, 'orig_failed_remaining': orig_files}

def get_output_and_backup__9fff5b0d3da3f15287b691b4b1944f6b_qw35sft2_e025a80d(env, config: dict):
    """Get content of output.txt and /tmp/output_backup.txt from the VM."""
    output_bytes = env.controller.get_file('/home/user/output.txt')
    backup_bytes = env.controller.get_file('/tmp/output_backup.txt')
    output_content = output_bytes.decode('utf-8', errors='replace') if output_bytes else None
    backup_content = backup_bytes.decode('utf-8', errors='replace') if backup_bytes else None
    return {'output_txt': output_content, 'backup_txt': backup_content}

def get_rename_and_move__eb200f79b699eb5d233b5379c7e1ea0e_qw35sft2_f80e9fe2(env, config: dict):
    """Check if todo_list_Jan_2 was renamed and moved to the home directory (not Desktop)."""
    vm_ip = env.vm_ip
    port = env.server_port
    r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/todo_list_Jan_2 && echo yes || echo no', 'shell': True})
    in_home = r1.status_code == 200 and r1.json().get('output', '').strip() == 'yes'
    r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'test -d ~/Desktop/todo_list_Jan_1 && echo yes || echo no', 'shell': True})
    old_gone = r2.status_code == 200 and r2.json().get('output', '').strip() == 'no'
    return {'in_home': in_home, 'old_gone': old_gone}

def get_volume_level__2ac92884cbfeefbd5a2d6929c20839d7_qw35sft2_ff2365b1(env, config: dict):
    """Get current PulseAudio sink volume output string."""
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pactl', 'get-sink-volume', '@DEFAULT_SINK@'], 'shell': False})
    if response.status_code == 200:
        return response.json().get('output', '')
    return ''

def get_restore_copy__f767d224a0cdea04b11256d9cffabd41_qw35sft2_5940f6b1(env, config: dict):
    """Check if poster is on Desktop (restored) and also copied to Documents folder."""
    vm_ip = env.vm_ip
    port = env.server_port
    try:
        resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': "echo 'DESKTOP:' && ls /home/user/Desktop/ 2>/dev/null && echo 'DOCUMENTS:' && ls /home/user/Documents/ 2>/dev/null", 'shell': True}, timeout=15)
        if resp.status_code != 200:
            return {'error': f'HTTP {resp.status_code}'}
        output = resp.json().get('output', '') or ''
        desktop_files = []
        docs_files = []
        section = None
        for line in output.splitlines():
            line = line.strip()
            if line == 'DESKTOP:':
                section = 'desktop'
            elif line == 'DOCUMENTS:':
                section = 'docs'
            elif line and section == 'desktop':
                desktop_files.append(line)
            elif line and section == 'docs':
                docs_files.append(line)
        return {'file_on_desktop': 'poster_party_night.webp' in desktop_files, 'file_in_documents': 'poster_party_night.webp' in docs_files}
    except Exception as e:
        logger_qw35sft2_fff6f2.error('get_restore_copy__f767d224a0cdea04b11256d9cffabd41 error: %s', e)
        return {'error': str(e)}

def get_copy_and_count_fails__8f62d3b2706e8f98a78e22aea189146a_qw35sft2_c7a85b79(env, config: dict):
    """Get *failed.ipynb files in ./fails and the content of ./fails/count.txt."""
    vm_ip = env.vm_ip
    port = env.server_port
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': "find /home/user/test_environment/fails -name '*failed.ipynb' -type f 2>/dev/null | sort", 'shell': True})
    fails_files = resp1.json().get('output', '').strip() if resp1.status_code == 200 else ''
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'cat /home/user/test_environment/fails/count.txt 2>/dev/null', 'shell': True})
    count_txt = resp2.json().get('output', '').strip() if resp2.status_code == 200 else ''
    return {'fails_files': fails_files, 'count_txt': count_txt}

def get_timezone_and_ntp__397a38a7528d157681d1e57316ff4073_qw35sft2_6e2cf5df(env, config: dict):
    """Get system timezone offset and NTP synchronization status.

    Returns a dict with:
      - 'timezone_offset': string like '+0000' from `date +%z`
      - 'ntp_status': string 'yes' or 'no' from `timedatectl show --property=NTP --value`
    """
    vm_ip = env.vm_ip
    port = env.server_port
    offset = ''
    ntp = ''
    try:
        r1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['date', '+%z'], 'shell': False}, timeout=10)
        if r1.status_code == 200:
            offset = r1.json().get('output', '').strip()
    except Exception as e:
        logger_qw35sft2_9ad635.warning(f'Failed to get timezone offset: {e}')
    try:
        r2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['timedatectl', 'show', '--property=NTP', '--value'], 'shell': False}, timeout=10)
        if r2.status_code == 200:
            ntp = r2.json().get('output', '').strip()
    except Exception as e:
        logger_qw35sft2_9ad635.warning(f'Failed to get NTP status: {e}')
    return {'timezone_offset': offset, 'ntp_status': ntp}

def get_volume_settings_state__1201dde5af3c06b515ed16eae4fad685_qw35sft2_6b004e3e(env, config: dict):
    """Get current volume level and whether gnome-control-center (Settings) is running."""
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    vol_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pactl', 'get-sink-volume', '@DEFAULT_SINK@'], 'shell': False})
    volume_output = ''
    if vol_resp.status_code == 200:
        volume_output = vol_resp.json().get('output', '')
    pgrep_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pgrep', '-x', 'gnome-control-center'], 'shell': False})
    settings_open = False
    if pgrep_resp.status_code == 200:
        pgrep_out = pgrep_resp.json().get('output', '').strip()
        settings_open = bool(pgrep_out)
    return {'volume_output': volume_output, 'settings_open': settings_open}

def get_gnome_favorites__f6775b09b9b85d12e5fafc15ee8143de_qw35sft2_cc93bd8a(env, config: dict):
    """Get the current GNOME shell favorite-apps list as a parsed Python list."""
    vm_ip = env.vm_ip
    port = env.server_port
    command = ['gsettings', 'get', 'org.gnome.shell', 'favorite-apps']
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
    if response.status_code == 200:
        apps_str = response.json()['output'].strip()
        try:
            apps = eval(apps_str)
            return {'apps': apps}
        except Exception:
            return {'error': f'Cannot parse favorites string: {apps_str}'}
    return {'error': 'Command execution failed'}

def get_power_triple_state__570bc607f27286d53f7b1a8eaf2af624_qw35sft2_a3f7abdb(env, config: dict):
    """Get idle-dim, idle-delay, and sleep-inactive-ac-type from GNOME power settings."""
    r = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power idle-dim', timeout=10)
    idle_dim_raw = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else str(r)).strip()
    r = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    idle_delay_raw = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else str(r)).strip()
    r = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power sleep-inactive-ac-type', timeout=10)
    suspend_type_raw = (r.get('output', r.get('stdout', '')) if isinstance(r, dict) else str(r)).strip()
    parts = idle_delay_raw.split()
    idle_delay = int(parts[-1]) if parts else -1
    return {'idle_dim': idle_dim_raw.lower().strip(), 'idle_delay': idle_delay, 'sleep_inactive_ac_type': suspend_type_raw.strip().strip("'")}

def get_notif_sounds__139fa3103c161a5284bcac4b5dc4f632_qw35sft2_a18ce23b(env, config: dict):
    """Get DND (show-banners) and system event-sounds settings."""
    r1 = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-banners', timeout=15)
    show_banners = (r1.get('output', r1.get('stdout', '')) if isinstance(r1, dict) else str(r1)).strip()
    r2 = env.controller.run_bash_script('gsettings get org.gnome.desktop.sound event-sounds', timeout=15)
    event_sounds = (r2.get('output', r2.get('stdout', '')) if isinstance(r2, dict) else str(r2)).strip()
    return {'show_banners': show_banners, 'event_sounds': event_sounds}

def get_gnome_favorites__fafd25459d3a377f8d0c6d91b93474de_qw35sft2_5897609a(env, config: dict):
    """Get the current GNOME shell favorite-apps list as a parsed Python list."""
    vm_ip = env.vm_ip
    port = env.server_port
    command = ['gsettings', 'get', 'org.gnome.shell', 'favorite-apps']
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
    if response.status_code == 200:
        apps_str = response.json()['output'].strip()
        try:
            apps = eval(apps_str)
            return {'apps': apps}
        except Exception:
            return {'error': f'Cannot parse favorites string: {apps_str}'}
    return {'error': 'Command execution failed'}

def get_power_dim_and_suspend__6cd1c7feeea3fc52512c1552e8d5b3b2_qw35sft2_dddf5aff(env, config: dict):
    """Get idle-dim and sleep-inactive-ac-type from GNOME power settings."""
    r1 = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power idle-dim', timeout=10)
    if isinstance(r1, dict):
        idle_dim_raw = r1.get('output', r1.get('stdout', '')).strip()
    else:
        idle_dim_raw = str(r1).strip()
    r2 = env.controller.run_bash_script('gsettings get org.gnome.settings-daemon.plugins.power sleep-inactive-ac-type', timeout=10)
    if isinstance(r2, dict):
        suspend_type_raw = r2.get('output', r2.get('stdout', '')).strip()
    else:
        suspend_type_raw = str(r2).strip()
    return {'idle_dim': idle_dim_raw.lower().strip(), 'sleep_inactive_ac_type': suspend_type_raw.strip().strip("'")}

def get_copy_and_rename__44eff79b79182e6eeb0c7debb1bc9f8c_qw35sft2_de77ce1d(env, config: dict):
    """Check file1 in dir1/dir2/dir3 and file1.bak in home dir."""
    result = env.controller.run_bash_script('echo dir1:$(test -f /home/user/dir1/file1 && echo yes || echo no)\necho dir2:$(test -f /home/user/dir2/file1 && echo yes || echo no)\necho dir3:$(test -f /home/user/dir3/file1 && echo yes || echo no)\necho bak:$(test -f /home/user/file1.bak && echo yes || echo no)', timeout=10)
    if not result:
        return {'error': 'command failed'}
    output = result.get('output', '')
    return {'dir1_has_file1': 'dir1:yes' in output, 'dir2_has_file1': 'dir2:yes' in output, 'dir3_has_file1': 'dir3:yes' in output, 'home_has_file1_bak': 'bak:yes' in output}

def get_accessibility_two_toggles__8ee8727a5fa2032ce96fa89a79ac9fab_qw35sft2_2815a142(env, config: dict):
    """Get large-text and high-contrast accessibility toggle states via gsettings."""
    vm_ip = env.vm_ip
    port = env.server_port
    result = {}
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.interface', 'large-text'], 'shell': False})
    if resp1.status_code == 200:
        result['large_text'] = resp1.json()['output'].strip() == 'true'
    else:
        result['large_text'] = None
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.interface', 'high-contrast'], 'shell': False})
    if resp2.status_code == 200:
        result['high_contrast'] = resp2.json()['output'].strip() == 'true'
    else:
        result['high_contrast'] = None
    return result

def get_eml_backup_state__d81f9d153146f886b82c4131d09d1ccb_qw35sft2_e1e59461(env, config: dict):
    """Return a dict with EML file count and directory listing for the backup directory."""
    directory = config.get('path', '/home/user/emails.bak/')
    vm_ip = env.vm_ip
    port = env.server_port
    count = '0'
    listing = ''
    try:
        count_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': f"find {directory} -maxdepth 1 -name '*.eml' 2>/dev/null | wc -l", 'shell': True})
        if count_resp.status_code == 200:
            count = count_resp.json().get('output', '0').strip()
    except Exception as e:
        logger_qw35sft2_edeb6d.error('Error counting EML files: %s', e)
    try:
        list_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['ls', directory], 'shell': False})
        if list_resp.status_code == 200:
            listing = list_resp.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_edeb6d.error('Error listing directory: %s', e)
    return {'count': count, 'listing': listing}

def get_bills_flag_and_tags__ed9097e7c75920408b4536024b8da2a4_qw35sft2_80b3d527(env, config: dict):
    """Read Bills mbox and return per-message starred status and Important tag ($label1)."""
    mbox_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills'
    file_bytes = env.controller.get_file(mbox_path)
    if not file_bytes:
        return {'error': 'Cannot read Bills mbox', 'message_count': 0, 'all_starred': False, 'all_important': False}
    content = file_bytes.decode('utf-8', errors='replace')
    messages = []
    parts = re.split('^From - ', content, flags=re.MULTILINE)
    for part in parts:
        if not part.strip():
            continue
        status_match = re.search('^X-Mozilla-Status: ([0-9A-Fa-f]+)', part, re.MULTILINE)
        if not status_match:
            continue
        status_int = int(status_match.group(1), 16)
        if status_int & 8:
            continue
        keys_match = re.search('^X-Mozilla-Keys:\\s*(.*?)$', part, re.MULTILINE)
        keys = keys_match.group(1).strip() if keys_match else ''
        messages.append({'is_starred': bool(status_int & 4), 'is_important': '$label1' in keys})
    return {'message_count': len(messages), 'all_starred': bool(messages) and all((m['is_starred'] for m in messages)), 'all_important': bool(messages) and all((m['is_important'] for m in messages))}

def get_bills_full_state__d21b0cc63afcca51c02fdd16b5d862e0_qw35sft2_0be511a0(env, config: dict):
    """Read Bills mbox and return starred, unread, and Important tag status for all messages."""
    mbox_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills'
    file_bytes = env.controller.get_file(mbox_path)
    if not file_bytes:
        return {'error': 'Cannot read Bills mbox', 'message_count': 0, 'all_starred': False, 'all_unread': False, 'all_important': False}
    content = file_bytes.decode('utf-8', errors='replace')
    messages = []
    parts = re.split('^From - ', content, flags=re.MULTILINE)
    for part in parts:
        if not part.strip():
            continue
        status_match = re.search('^X-Mozilla-Status: ([0-9A-Fa-f]+)', part, re.MULTILINE)
        if not status_match:
            continue
        status_int = int(status_match.group(1), 16)
        if status_int & 8:
            continue
        keys_match = re.search('^X-Mozilla-Keys:\\s*(.*?)$', part, re.MULTILINE)
        keys = keys_match.group(1).strip() if keys_match else ''
        messages.append({'is_starred': bool(status_int & 4), 'is_read': bool(status_int & 1), 'is_important': '$label1' in keys})
    return {'message_count': len(messages), 'all_starred': bool(messages) and all((m['is_starred'] for m in messages)), 'all_unread': bool(messages) and all((not m['is_read'] for m in messages)), 'all_important': bool(messages) and all((m['is_important'] for m in messages))}

def get_eml_count__2731b9abd5cfbad9ed4df8aae737addc_qw35sft2_d185deda(env, config: dict):
    """Count .eml files in the specified directory on the VM."""
    directory = config.get('path', '/home/user/emails.bak/')
    vm_ip = env.vm_ip
    port = env.server_port
    command = f"find {directory} -maxdepth 1 -name '*.eml' 2>/dev/null | wc -l"
    try:
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
        if response.status_code == 200:
            return response.json().get('output', '0').strip()
        logger_qw35sft2_f39cbb.error('Failed to execute find command, status: %d', response.status_code)
        return '0'
    except Exception as e:
        logger_qw35sft2_f39cbb.error('Error counting EML files: %s', e)
        return '0'

def get_bills_flag_state__8e8eeb1588f1109e98ccb02a0faa787c_qw35sft2_87acafa7(env, config: dict):
    """Read Bills mbox and return per-message starred/unread status."""
    mbox_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills'
    file_bytes = env.controller.get_file(mbox_path)
    if not file_bytes:
        return {'error': 'Cannot read Bills mbox', 'message_count': 0, 'all_starred': False, 'all_unread': False}
    content = file_bytes.decode('utf-8', errors='replace')
    messages = []
    parts = re.split('^From - ', content, flags=re.MULTILINE)
    for part in parts:
        if not part.strip():
            continue
        status_match = re.search('^X-Mozilla-Status: ([0-9A-Fa-f]+)', part, re.MULTILINE)
        if not status_match:
            continue
        status_int = int(status_match.group(1), 16)
        if status_int & 8:
            continue
        messages.append({'is_starred': bool(status_int & 4), 'is_read': bool(status_int & 1)})
    return {'message_count': len(messages), 'all_starred': bool(messages) and all((m['is_starred'] for m in messages)), 'all_unread': bool(messages) and all((not m['is_read'] for m in messages))}

def get_eml_listing__dfeb48225188ee18fb4de9d6f0048829_qw35sft2_71b752de(env, config: dict):
    """List files in the specified backup directory on the VM."""
    directory = config.get('path', '/home/user/emails.bak/')
    vm_ip = env.vm_ip
    port = env.server_port
    try:
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['ls', directory], 'shell': False})
        if response.status_code == 200:
            return response.json().get('output', '')
        logger_qw35sft2_0d8aba.error('Failed to list directory, status: %d', response.status_code)
        return ''
    except Exception as e:
        logger_qw35sft2_0d8aba.error('Error listing EML files: %s', e)
        return ''

def get_bills_and_filter__44c5673d772aeaf1fd4e6dcc8a32111d_qw35sft2_a2ee062f(env, config: dict):
    """Read Bills mbox starred status AND check for a 'Mark flagged' message filter."""
    mbox_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills'
    file_bytes = env.controller.get_file(mbox_path)
    all_starred = False
    message_count = 0
    if file_bytes:
        content = file_bytes.decode('utf-8', errors='replace')
        messages = []
        parts = re.split('^From - ', content, flags=re.MULTILINE)
        for part in parts:
            if not part.strip():
                continue
            status_match = re.search('^X-Mozilla-Status: ([0-9A-Fa-f]+)', part, re.MULTILINE)
            if not status_match:
                continue
            status_int = int(status_match.group(1), 16)
            if status_int & 8:
                continue
            messages.append(bool(status_int & 4))
        message_count = len(messages)
        all_starred = bool(messages) and all(messages)
    filter_path = '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/msgFilterRules.dat'
    filter_bytes = env.controller.get_file(filter_path)
    filter_has_star_action = False
    if filter_bytes:
        filter_content = filter_bytes.decode('utf-8', errors='replace')
        blocks = re.split('\\n\\s*\\n', filter_content)
        for block in blocks:
            if 'Mark flagged' in block and 'Invoice' in block:
                filter_has_star_action = True
                break
    return {'message_count': message_count, 'all_starred': all_starred, 'filter_has_star_action': filter_has_star_action}

def get_ext_and_multi_settings__afc2fb3b68b53df8d476e05f07933aff_qw35sft2_71352dcc(env, config: dict):
    """Get extension list, editor.tabSize, and editor.formatOnSave from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_fa1173.warning('Failed to get extension list: %s', e)
    tab_size = None
    format_on_save = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            tab_size = settings.get('editor.tabSize')
            format_on_save = settings.get('editor.formatOnSave')
    except Exception as e:
        logger_qw35sft2_fa1173.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'tab_size': tab_size, 'format_on_save': format_on_save}

def get_ext_and_fontsize__fe14a817663aeade8a20cb0f5baad2b6_qw35sft2_e0fd936f(env, config: dict):
    """Get extension list and editor.fontSize from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_c93c96.warning('Failed to get extension list: %s', e)
    font_size = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            font_size = settings.get('editor.fontSize')
    except Exception as e:
        logger_qw35sft2_c93c96.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'font_size': font_size}

def get_ext_and_wordwrap__324a83eafb9ff6ed07fafe7e199af4d5_qw35sft2_1f593828(env, config: dict):
    """Get extension list and editor.wordWrap from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_2620a4.warning('Failed to get extension list: %s', e)
    word_wrap = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            word_wrap = settings.get('editor.wordWrap')
    except Exception as e:
        logger_qw35sft2_2620a4.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'word_wrap': word_wrap}

def get_ext_and_settings__bb23bc65a82d76a52cba97bb0a6bf9bc_qw35sft2_923dbd2e(env, config: dict):
    """Get extension list and editor.fontSize from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_5d3c04.warning('Failed to get extension list: %s', e)
    font_size = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            font_size = settings.get('editor.fontSize')
    except Exception as e:
        logger_qw35sft2_5d3c04.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'font_size': font_size}

def get_dual_ext__95ee9f441436911870f312520ed4f195_qw35sft2_08027149(env, config: dict):
    """Get installed VS Code extension list to verify two extensions."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_2ed96e.warning('Failed to get extension list: %s', e)
    return {'ext_list': ext_list}

def get_ext_and_wordwrap__6ac7db98d9670abbcf37d963dc27bc84_qw35sft2_f286792b(env, config: dict):
    """Get extension list and editor.wordWrap from settings.json."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_b0d92d.warning('Failed to get extension list: %s', e)
    word_wrap = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            word_wrap = settings.get('editor.wordWrap')
    except Exception as e:
        logger_qw35sft2_b0d92d.warning('Failed to read settings.json: %s', e)
    return {'ext_list': ext_list, 'word_wrap': word_wrap}
