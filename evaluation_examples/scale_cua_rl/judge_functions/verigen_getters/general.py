"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_python_syntax_check__94eb7a158102fbe4b2865eae8fd2e0f7', 'get_git_push_and_branch__216e252ae5e5ecc8fbd5efc80755d934_qw35sft2_e7b5c26c', 'get_python_docs_and_font__492670c8dd0e3374d756c80fee290bb9_qw35sft2_4227e830', 'get_python_env_setup__be89f5274d64cd2fc7ba060d990658f8_qw35sft2_70da97e9', 'get_bashrc_python_settings__dd32b078ecc7acd064c9c431fc6fece6_qw35sft2_eedcd129', 'get_volume_terminal_state__bafe7fff6ff48b4d9b5b1fd5f2d7f539_qw35sft2_83d48ccf', 'get_os_night_dark_state__0de6f297debf99316b6fa7b7c038fa5d_qw35sft2_a81b62ba']

def get_python_syntax_check__94eb7a158102fbe4b2865eae8fd2e0f7(env, config: dict):
    """Download Python file and check if it has valid syntax."""
    import ast
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'valid': False, 'error': 'File not found', 'content': ''}
    content = file_bytes.decode('utf-8', errors='replace')
    try:
        ast.parse(content)
        return {'valid': True, 'content': content}
    except SyntaxError as e:
        return {'valid': False, 'error': str(e), 'content': content}

def get_git_push_and_branch__216e252ae5e5ecc8fbd5efc80755d934_qw35sft2_e7b5c26c(env, config: dict):
    """
    Get the latest commit message in the remote repo and the list of local branches
    in the binder project.
    Returns dict with keys 'remote_log' and 'branches'.
    """
    remote_log_result = env.controller.run_bash_script('git -C /home/user/projects/remote_project log --oneline -1 2>&1 || echo "GIT_ERROR"', timeout=30)
    if isinstance(remote_log_result, dict):
        remote_log = remote_log_result.get('stdout', '') or remote_log_result.get('output', '') or str(remote_log_result)
    else:
        remote_log = str(remote_log_result) if remote_log_result else ''
    branches_result = env.controller.run_bash_script('git -C /home/user/projects/binder branch 2>&1', timeout=30)
    if isinstance(branches_result, dict):
        branches = branches_result.get('stdout', '') or branches_result.get('output', '') or str(branches_result)
    else:
        branches = str(branches_result) if branches_result else ''
    return {'remote_log': remote_log.strip(), 'branches': branches.strip()}

def get_python_docs_and_font__492670c8dd0e3374d756c80fee290bb9_qw35sft2_4227e830(env, config: dict):
    """Get active Chrome tab URL and current default font size from Chrome preferences."""
    import json
    font_size = None
    try:
        prefs_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Preferences')
        if prefs_bytes:
            prefs = json.loads(prefs_bytes.decode('utf-8'))
            font_size = prefs.get('webkit', {}).get('webprefs', {}).get('default_font_size')
    except Exception:
        pass
    active_url = ''
    try:
        tree = env.controller.get_accessibility_tree()
        if tree:
            import re
            match = re.search('Address and search bar[^\\n]*\\nvalue: ([^\\n]+)', tree)
            if match:
                active_url = match.group(1).strip()
            else:
                match = re.search('https?://[^\\s\\\'"<>]+', tree)
                if match:
                    active_url = match.group(0).strip()
    except Exception:
        pass
    return {'active_url': active_url, 'font_size': font_size}

def get_python_env_setup__be89f5274d64cd2fc7ba060d990658f8_qw35sft2_70da97e9(env, config: dict):
    """
    Read /etc/environment and ~/python_info.txt for a two-goal Python config check.
    """
    vm_ip = env.vm_ip
    port = env.server_port
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['cat', '/etc/environment'], 'shell': False})
    env_content = resp1.json().get('output', '') if resp1.status_code == 200 else ''
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['cat', '/home/user/python_info.txt'], 'shell': False})
    file_content = resp2.json().get('output', '') if resp2.status_code == 200 else ''
    return {'env_content': env_content, 'file_content': file_content}

def get_bashrc_python_settings__dd32b078ecc7acd064c9c431fc6fece6_qw35sft2_eedcd129(env, config: dict):
    """Read ~/.bashrc content for Python-related alias and env var checks."""
    vm_ip = env.vm_ip
    port = env.server_port
    resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['cat', '/home/user/.bashrc'], 'shell': False})
    if resp.status_code == 200:
        content = resp.json().get('output', '')
    else:
        content = ''
    return {'bashrc_content': content}

def get_volume_terminal_state__bafe7fff6ff48b4d9b5b1fd5f2d7f539_qw35sft2_83d48ccf(env, config: dict):
    """Get current volume level and whether gnome-terminal is running."""
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    vol_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pactl', 'get-sink-volume', '@DEFAULT_SINK@'], 'shell': False})
    volume_output = ''
    if vol_resp.status_code == 200:
        volume_output = vol_resp.json().get('output', '')
    pgrep_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['pgrep', 'gnome-terminal'], 'shell': False})
    terminal_open = False
    if pgrep_resp.status_code == 200:
        pgrep_out = pgrep_resp.json().get('output', '').strip()
        terminal_open = bool(pgrep_out)
    return {'volume_output': volume_output, 'terminal_open': terminal_open}

def get_os_night_dark_state__0de6f297debf99316b6fa7b7c038fa5d_qw35sft2_a81b62ba(env, config: dict):
    """Get Night Light enabled state and color-scheme from GNOME gsettings."""
    night_result = env.controller.run_command(['gsettings', 'get', 'org.gnome.settings-daemon.plugins.color', 'night-light-enabled'])
    dark_result = env.controller.run_command(['gsettings', 'get', 'org.gnome.desktop.interface', 'color-scheme'])
    night_output = (night_result or {}).get('output', '').strip()
    dark_output = (dark_result or {}).get('output', '').strip()
    return {'night_light_enabled': night_output, 'color_scheme': dark_output}
