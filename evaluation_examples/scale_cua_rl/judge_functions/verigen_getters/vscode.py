"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_vscode_window_title__2e0260d8e24ca4d49ed462b4c3422ce2', 'get_vscode_gitlens_wordwrap__5c2597160999bc0ba0a4a23d777e3ec1', 'get_vscode_settings__215b1cc2c3da9fb70c240f9f4f8df6b5', 'get_vscode_file_content__43d92268ad8323e8d3f0af9cc0d7b2ce', 'get_vscode_settings__ce8e3f26fe80873188ba40e45449d3b9', 'get_workspace_settings__6a6725db0c0fb336a2ee8d61d5b62877', 'get_vscode_file_content__7061bd4dde7607bb250f2ba7725dc0d5', 'get_vscode_settings_json__24dcbea28334e3c805df8afddc88a5f1', 'get_vscode_settings__ebe9d7bea41ec6a21d7329851f18b701', 'get_vscode_extensions__c324ab1edeaa19edae46a63b788543c8', 'get_vscode_settings__b082db988cea152fff3c88fce0623d97', 'get_workspace_settings__50a1f916fbc8eceedebd80aa7a09697c', 'get_vscode_settings__ef57df2dacd824ebd85cb86fda295eab', 'get_vscode_file_lines__c8cb9d1d5c3b002c8ddbeb93d3580891', 'get_workspace_file_content__56abb1eed13e878c335e4f175376e020', 'get_vscode_settings__c7225cce9efdba9aa8af211dee6285ef', 'get_vscode_file_lines__c54a7288908815335d3763c41c9cef77', 'get_vscode_file_lines__350a72c316de89189f7f590d79389d9e', 'get_vscode_file_content__12dc12177d981fb82b80cf3f9509e0cf', 'get_workspace_full_state__218f219ceca7e7e0003dc76221985da3_qw35sft2_aa8dbe95', 'get_vscode_debug_focus_console__b4f4f2d972bad7dfc4a5dbad64bc0fb1_qw35sft2_5e672e89', 'get_vscode_theme_font_wrap__12b042cae58f794626351d3c568391b7_qw35sft2_9b58fcc5', 'get_vscode_keybindings__ed42af8e24d428ca40b99f135c68274f_qw35sft2_44312464', 'get_vscode_settings_dict__81417cbab160b1f4e52b4e40c1e83805_qw35sft2_cdf64786', 'get_vscode_settings__e20226d1d49056af54e3e69873c8fc3c_qw35sft2_46057785', 'get_vscode_settings__193e70599e6295621968803241fc32b6_qw35sft2_4691c4e5', 'get_vscode_word_wrap__15748b95140e1ac4229209dbc7dfe067_qw35sft2_c25b0e86', 'get_vscode_remove_two_list_shortcuts__45d3929cc132a565457f8b673ec620af_qw35sft2_d2938a16', 'get_workspace_and_file__5eb515badd2fc6ca56c5f84fbe245437_qw35sft2_5488a104', 'get_vscode_wrap_and_tab__5464614ee51132440430b0764a874f9a_qw35sft2_59244c58', 'get_vscode_txt_content__ae0b3ab6740bff99e232fccee3b36aad_qw35sft2_c0eb958d', 'get_vscode_exclude_dual__26878716b564c56320e14f3b9550ebe4_qw35sft2_1c113466', 'get_vscode_ext_list__eb3392bef87727a9af1f8c1971adb226_qw35sft2_1720c58f', 'get_vscode_file_and_settings__9fe429b95a40aa1badcb7de0b8a6b0f1_qw35sft2_3757270c', 'get_vscode_locale_and_settings__de00e991b11dacec5fa1917014fc1ed1_qw35sft2_973307ec', 'get_workspace_and_fontsize__4cb0241f3e1c4176b65966dc96627af4_qw35sft2_73075398', 'get_vscode_py_indent_docstring__3d465870a880c6e96493a5637b57636e_qw35sft2_e21fc5b4', 'get_vscode_debug_focus_font__802820e3fd86caf9b92a9751803125a2_qw35sft2_77839c8f', 'get_vscode_settings_dict__8bfa6a8d7119af313a0660a3dea4ceab_qw35sft2_13b525eb', 'get_vscode_keybindings__70885a60e67c4def20a9d0e3284c1793_qw35sft2_d6793f48', 'get_vscode_settings__fa93c32e8afd83b84bab5f585818b0fb_qw35sft2_e1ecdb7a', 'get_ext_and_keybindings__44f31fc07357666acac95e9a419ee84f_qw35sft2_560113fa', 'get_vscode_theme_wordwrap__978e7cba9602912bc979b36fdc172116_qw35sft2_6665e744', 'get_vscode_txt_content__004ad4b9f1b7c0d57d366b67195586e3_qw35sft2_8fb86838', 'get_workspace_folders__f72e9e4f450564b5101c9194862d205b_qw35sft2_32d45db6', 'get_vscode_settings__d45b70bac3694dd736206b10dff26100_qw35sft2_cf1fa267', 'get_vscode_exclude_pytest__6a66711a30fb269a0e7d77ca8b497959_qw35sft2_2762ee26', 'get_vscode_remove_and_add_keybinding__8bf5980383e9d80a451eac2b89ccf1df_qw35sft2_2666838f', 'get_vscode_word_wrap_settings__35c11665b6e5c9d9c9abe7f855ccfe46_qw35sft2_80c5a18c', 'get_vscode_locale_and_settings__e258fe9bb867a78f6ed6748a6bee2270_qw35sft2_d3d959ab', 'get_vscode_ext_and_settings__e2317900e4eb6fe8b0284d167aa38270_qw35sft2_2f9de2bd', 'get_workspace_and_extension__bead7b293cda299d38413ae3ff873cc8_qw35sft2_9fcd0186', 'get_vscode_py_function_file__9e0f9c32752df4ad0956b05cabebbc98_qw35sft2_e6aaeebb', 'get_vscode_workspace_in_storage__5752e6d2cd4f13c8d24f076655028c23_qw35sft2_a6c9eb8a', 'get_vscode_py_indent_blankline__2c836d0eca9976b79042c49c51fe71cf_qw35sft2_67fb627b', 'get_vscode_settings_dict__9ec2182c98e7d1d45bbbb3ae2d0db42e_qw35sft2_be433eca', 'get_vscode_debug_focus_wrap__09cc1b24020e8f7c0f7b48c6a4d96413_qw35sft2_4aeffc35', 'get_vscode_settings__a4b8eefcd32d2ac2b259c28c71814550_qw35sft2_af664dff', 'get_workspace_folders__d5491cdef8176b887dce317a57139150_qw35sft2_81e65fe6', 'get_vscode_settings__5cd158cc213aefb25242b372a3862ff3_qw35sft2_3d23d8ed', 'get_vscode_txt_content__076bfd2d3ed2cd7e42a937a1dd80c8a9_qw35sft2_b2456901', 'get_vscode_wrap_and_format__95a613ea8825af09f40057c978f02ad7_qw35sft2_f6281378', 'get_workspace_and_wordwrap__1713b5c17b096bf267a520ed66d388f0_qw35sft2_f0f8a503', 'get_vscode_exclude_autosave__7797789b3c2086d41b9d5f1699195599_qw35sft2_4b32b4aa', 'get_vscode_desktop_file_comment__eea98fddd6e0176c917193edbee2b203_qw35sft2_3e7e2657', 'get_vscode_ext_and_settings__88ac2a64a7ee030ff707715de339ea34_qw35sft2_c88a3e48', 'get_vscode_settings_dict__d122d575a22853e47d56ee92b6c18378_qw35sft2_b8e62d25', 'get_vscode_txt_content__e470774862e9eb18ed4fb0ad458ab39f_qw35sft2_abefc668', 'get_vscode_open_workspace__c37846152daf06e05c18277c5b1f636d_qw35sft2_5b8c70c2', 'get_workspace_folders__edd7ecbec0169642eef724f574a58652_qw35sft2_04296562', 'get_vscode_workspace_state__f91feeeb0a889e0b6b274960d3d06c7d_qw35sft2_cc26b79e', 'get_vscode_tab_size__015bdcf4a8f96215b5579fc53db2d526_qw35sft2_860be523', 'get_vscode_debug_focus_terminal__29baecbd3036f18eff818341b1c0e4e8_qw35sft2_0d7e8ce8', 'get_vscode_exclude_fontsize__33cb4b35f94afe16fa1fe71880064b0b_qw35sft2_745131bd', 'get_vscode_settings__d97bffb72b4b8c27def4b087597b01a0_qw35sft2_5da34fe4', 'get_vscode_ext_list__521eb992bf97de2d79e4aa6c081db684_qw35sft2_67841176', 'get_vscode_txt_content__d87d9b6cb6d6a7fc9d5d36ba0c96e0de_qw35sft2_c1aa3798', 'get_vscode_active_file__6a582969bc6d0b9624887969279510e3_qw35sft2_2cbb6469', 'get_vscode_settings_dict__7e25498cca269e458737cc32ebf8beb0_qw35sft2_477d89cd']

def get_vscode_window_title__2e0260d8e24ca4d49ed462b4c3422ce2(env, config: dict):
    """Get all VS Code window titles from the desktop."""
    result = env.controller.run_bash_script("wmctrl -l 2>/dev/null | grep -i 'Visual Studio Code' || xdotool search --name 'Visual Studio Code' getwindowname 2>/dev/null", timeout=10)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '')
    elif isinstance(result, str):
        output = result
    return {'window_titles': output}

def get_vscode_gitlens_wordwrap__5c2597160999bc0ba0a4a23d777e3ec1(env, config: dict):
    """Get GitLens extension installation status and editor.wordWrap setting."""
    result = {}
    ext_output = env.controller.run_bash_script('code --list-extensions 2>/dev/null', timeout=30)
    ext_list = ext_output.get('output', '') if isinstance(ext_output, dict) else str(ext_output)
    result['extensions'] = ext_list
    settings_path = '/home/user/.config/Code/User/settings.json'
    try:
        file_bytes = env.controller.get_file(settings_path)
        if file_bytes:
            settings = json.loads(file_bytes.decode('utf-8'))
            result['settings'] = settings
        else:
            result['settings'] = {}
    except Exception:
        result['settings'] = {}
    return result

def get_vscode_settings__215b1cc2c3da9fb70c240f9f4f8df6b5(env, config: dict):
    """Get VS Code user settings from settings.json."""
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_file_content__43d92268ad8323e8d3f0af9cc0d7b2ce(env, config: dict):
    """Read text file content from VM."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': str(e), 'content': ''}
    return {'content': content}

def get_vscode_settings__ce8e3f26fe80873188ba40e45449d3b9(env, config: dict):
    """Get VS Code user settings.json content."""
    settings_path = config.get('path', '/home/user/.config/Code/User/settings.json')
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_workspace_settings__6a6725db0c0fb336a2ee8d61d5b62877(env, config: dict):
    """Get workspace settings.json content as parsed JSON."""
    import json
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        content = json.loads(file_bytes.decode('utf-8'))
        return {'settings': content}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_file_content__7061bd4dde7607bb250f2ba7725dc0d5(env, config: dict):
    """Read text file content from VM."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': str(e), 'content': ''}
    return {'content': content}

def get_vscode_settings_json__24dcbea28334e3c805df8afddc88a5f1(env, config: dict):
    """Read VS Code settings.json from VM."""
    import json
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
    except Exception:
        return {'error': 'Failed to read settings file'}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
        settings = json.loads(content)
        return {'settings': settings}
    except json.JSONDecodeError:
        return {'error': 'Invalid JSON in settings file'}

def get_vscode_settings__ebe9d7bea41ec6a21d7329851f18b701(env, config: dict):
    """Get VS Code user settings.json content."""
    settings_path = config.get('path', '/home/user/.config/Code/User/settings.json')
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_extensions__c324ab1edeaa19edae46a63b788543c8(env, config: dict):
    """Get list of installed VS Code extensions."""
    result = env.controller.run_bash_script('code --list-extensions', timeout=30)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', '')
    elif isinstance(result, str):
        output = result
    return {'extensions': output}

def get_vscode_settings__b082db988cea152fff3c88fce0623d97(env, config: dict):
    """Get VS Code user settings.json content."""
    settings_path = config.get('path', '/home/user/.config/Code/User/settings.json')
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_workspace_settings__50a1f916fbc8eceedebd80aa7a09697c(env, config: dict):
    """Read a VS Code workspace file and extract the settings object."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        settings = data.get('settings', {})
        return {'settings': settings}
    except (json.JSONDecodeError, UnicodeDecodeError) as e:
        return {'error': f'Failed to parse workspace file: {str(e)}'}

def get_vscode_settings__ef57df2dacd824ebd85cb86fda295eab(env, config: dict):
    """Get VS Code user settings from settings.json."""
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_file_lines__c8cb9d1d5c3b002c8ddbeb93d3580891(env, config: dict):
    """Read file from VM and return all lines."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = content.split('\n')
    if lines and lines[-1] == '':
        lines = lines[:-1]
    return {'lines': lines, 'num_lines': len(lines)}

def get_workspace_file_content__56abb1eed13e878c335e4f175376e020(env, config: dict):
    """Read a .code-workspace file from the VM and parse its folder structure."""
    import json
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'folders': []}
    try:
        content = json.loads(file_bytes.decode('utf-8'))
        folders = [f.get('path', '') for f in content.get('folders', [])]
        return {'folders': folders, 'folder_count': len(folders)}
    except Exception as e:
        return {'error': str(e), 'folders': []}

def get_vscode_settings__c7225cce9efdba9aa8af211dee6285ef(env, config: dict):
    """Get VS Code user settings from settings.json."""
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'settings': settings}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_file_lines__c54a7288908815335d3763c41c9cef77(env, config: dict):
    """Read file from VM and return all lines."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = content.split('\n')
    if lines and lines[-1] == '':
        lines = lines[:-1]
    return {'lines': lines, 'num_lines': len(lines)}

def get_vscode_file_lines__350a72c316de89189f7f590d79389d9e(env, config: dict):
    """Read file from VM and return all lines."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    lines = content.split('\n')
    if lines and lines[-1] == '':
        lines = lines[:-1]
    return {'lines': lines, 'num_lines': len(lines)}

def get_vscode_file_content__12dc12177d981fb82b80cf3f9509e0cf(env, config: dict):
    """Read text file content from VM."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': str(e), 'content': ''}
    return {'content': content}

def get_workspace_full_state__218f219ceca7e7e0003dc76221985da3_qw35sft2_aa8dbe95(env, config: dict):
    """Get full workspace state: terminal CWD, open Chrome tab URLs, and bookmarks.

    Returns dict with:
    - 'terminal_cwd': str
    - 'open_urls': list of str
    - 'bookmarks': nested bookmark dict
    """
    import json
    terminal_cwd = ''
    try:
        result = env.controller.run_bash_script('for pid in $(pgrep -x bash 2>/dev/null); do cwd=$(readlink /proc/$pid/cwd 2>/dev/null); if [ -n "$cwd" ]; then echo "$cwd"; fi; done | head -5', timeout=15)
        if result and isinstance(result, dict):
            output = result.get('output', '') or result.get('stdout', '')
        elif isinstance(result, str):
            output = result
        else:
            output = ''
        lines = [l.strip() for l in output.strip().splitlines() if l.strip()]
        if lines:
            terminal_cwd = lines[-1]
    except Exception:
        pass
    open_urls = []
    try:
        tab_result = env.controller.run_bash_script('python3 -c "import urllib.request, json; data=urllib.request.urlopen(\'http://localhost:9222/json\').read(); tabs=json.loads(data); print(json.dumps([t.get(\'url\',\'\') for t in tabs if t.get(\'type\')==\'page\']))" 2>/dev/null || echo \'[]\'', timeout=15)
        if tab_result and isinstance(tab_result, dict):
            output = tab_result.get('output', '') or tab_result.get('stdout', '')
        elif isinstance(tab_result, str):
            output = tab_result
        else:
            output = '[]'
        output = output.strip()
        if output:
            open_urls = json.loads(output)
    except Exception:
        pass
    bookmarks = {}
    try:
        bm_bytes = env.controller.get_file('/home/user/.config/google-chrome/Default/Bookmarks')
        if bm_bytes:
            bookmarks = json.loads(bm_bytes.decode('utf-8'))
    except Exception:
        pass
    return {'terminal_cwd': terminal_cwd, 'open_urls': open_urls, 'bookmarks': bookmarks}

def get_vscode_debug_focus_console__b4f4f2d972bad7dfc4a5dbad64bc0fb1_qw35sft2_5e672e89(env, config: dict):
    """Get debug.focusEditorOnBreak and debug.internalConsoleOptions from settings.json."""
    try:
        file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings = json.loads(file_bytes.decode('utf-8'))
    except Exception:
        settings = {}
    return {'focus_editor_on_break': settings.get('debug.focusEditorOnBreak', True), 'internal_console_options': settings.get('debug.internalConsoleOptions', 'openOnFirstSessionStart')}

def get_vscode_theme_font_wrap__12b042cae58f794626351d3c568391b7_qw35sft2_9b58fcc5(env, config: dict):
    """Read VS Code settings.json and return colorTheme, editor.fontSize, and editor.wordWrap."""
    path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'colorTheme': settings.get('workbench.colorTheme'), 'fontSize': settings.get('editor.fontSize'), 'wordWrap': settings.get('editor.wordWrap')}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_keybindings__ed42af8e24d428ca40b99f135c68274f_qw35sft2_44312464(env, config: dict):
    """Read VS Code keybindings.json and return parsed entries."""
    keybindings_path = '/home/user/.config/Code/User/keybindings.json'
    file_bytes = env.controller.get_file(keybindings_path)
    if not file_bytes:
        return {'error': 'keybindings.json not found', 'entries': []}
    try:
        content = file_bytes.decode('utf-8')
        content = re.sub('//[^\\n]*', '', content)
        content = re.sub('/\\*.*?\\*/', '', content, flags=re.DOTALL)
        entries = json.loads(content.strip())
        return {'entries': entries if isinstance(entries, list) else []}
    except Exception as e:
        return {'error': str(e), 'entries': []}

def get_vscode_settings_dict__81417cbab160b1f4e52b4e40c1e83805_qw35sft2_cdf64786(env, config: dict):
    """Read VS Code user settings.json and return parsed dict."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_settings__e20226d1d49056af54e3e69873c8fc3c_qw35sft2_46057785(env, config: dict):
    """Read VS Code User settings.json from the VM and return its contents as a dict."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8')
        settings = json.loads(content)
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_settings__193e70599e6295621968803241fc32b6_qw35sft2_4691c4e5(env, config: dict):
    """Read VS Code user settings.json and return its parsed content."""
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8') if isinstance(file_bytes, bytes) else file_bytes
        return json.loads(content)
    except Exception as e:
        return {'error': str(e)}

def get_vscode_word_wrap__15748b95140e1ac4229209dbc7dfe067_qw35sft2_c25b0e86(env, config: dict):
    """Read editor.wordWrap from VS Code settings.json."""
    settings_path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        value = settings.get('editor.wordWrap')
        return {'value': value}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_remove_two_list_shortcuts__45d3929cc132a565457f8b673ec620af_qw35sft2_d2938a16(env, config: dict):
    """Read keybindings.json and check for both ctrl+f (list.find) and Escape (list.closeFind) removals."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/keybindings.json')
    if not file_bytes:
        return {'error': 'keybindings.json not found'}
    try:
        content = file_bytes.decode('utf-8').strip()
        if not content or content == '[]':
            return {'ctrlf_removed': False, 'escape_removed': False}
        bindings = json.loads(content)
        ctrlf_removed = False
        escape_removed = False
        for b in bindings:
            key = b.get('key', '').lower().replace(' ', '')
            cmd = b.get('command', '')
            when = b.get('when', '')
            if key == 'ctrl+f' and cmd == '-list.find' and ('listFocus' in when) and ('listSupportsFind' in when):
                ctrlf_removed = True
            if key == 'escape' and cmd == '-list.closeFind' and ('listFocus' in when):
                escape_removed = True
        return {'ctrlf_removed': ctrlf_removed, 'escape_removed': escape_removed}
    except Exception as e:
        return {'error': str(e)}

def get_workspace_and_file__5eb515badd2fc6ca56c5f84fbe245437_qw35sft2_5488a104(env, config: dict):
    """Read workspace folders and check whether a specific file exists in the VM."""
    workspace_path = config.get('path', '/home/user/project.code-workspace')
    check_file = config.get('check_file', '/home/user/data1/notes.txt')
    folder_names = []
    file_bytes = env.controller.get_file(workspace_path)
    if file_bytes:
        try:
            data = json.loads(file_bytes.decode('utf-8'))
            raw_paths = [f.get('path', '') for f in data.get('folders', [])]
            ws_dir = os.path.dirname(workspace_path)
            for p in raw_paths:
                if os.path.isabs(p):
                    folder_names.append(os.path.basename(p.rstrip('/')))
                else:
                    folder_names.append(os.path.basename(os.path.normpath(os.path.join(ws_dir, p))))
        except Exception:
            pass
    file_exists = False
    try:
        bash_result = env.controller.run_bash_script(f'test -f "{check_file}" && echo "exists" || echo "absent"', timeout=10)
        if isinstance(bash_result, dict):
            output = bash_result.get('output', '') or bash_result.get('stdout', '')
        else:
            output = str(bash_result)
        file_exists = 'exists' in output
    except Exception:
        pass
    return {'folder_names': folder_names, 'file_exists': file_exists}

def get_vscode_wrap_and_tab__5464614ee51132440430b0764a874f9a_qw35sft2_59244c58(env, config: dict):
    """Read VS Code user settings.json and return wordWrapColumn and tabSize values."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'wordWrapColumn': settings.get('editor.wordWrapColumn'), 'tabSize': settings.get('editor.tabSize')}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_txt_content__ae0b3ab6740bff99e232fccee3b36aad_qw35sft2_c0eb958d(env, config: dict):
    """Read vscode_replace_text.txt from the VM and return its text content."""
    path = config.get('path', '/home/user/Desktop/vscode_replace_text.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found or empty'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {str(e)}'}
    return {'content': content}

def get_vscode_exclude_dual__26878716b564c56320e14f3b9550ebe4_qw35sft2_1c113466(env, config: dict):
    """Get VS Code files.exclude dict from settings.json."""
    files_exclude = {}
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            files_exclude = settings.get('files.exclude', {})
    except Exception as e:
        logger_qw35sft2_4c210b.warning('Failed to read settings.json: %s', e)
    return {'files_exclude': files_exclude}

def get_vscode_ext_list__eb3392bef87727a9af1f8c1971adb226_qw35sft2_1720c58f(env, config: dict):
    """Get installed VS Code extension list."""
    result = env.controller.run_bash_script('code --list-extensions 2>/dev/null', timeout=30)
    if isinstance(result, dict):
        output = result.get('output', '') or result.get('stdout', '') or ''
    else:
        output = str(result) if result else ''
    return {'extensions': output.strip().lower()}

def get_vscode_file_and_settings__9fe429b95a40aa1badcb7de0b8a6b0f1_qw35sft2_3757270c(env, config: dict):
    """Get test.py existence and VS Code tab size setting."""
    import json
    file_path = config.get('path', '/home/user/Desktop/test.py')
    settings_path = config.get('settings_path', '/home/user/.config/Code/User/settings.json')
    file_exists = False
    try:
        data = env.controller.get_file(file_path)
        file_exists = data is not None
    except Exception:
        pass
    tab_size = None
    try:
        settings_bytes = env.controller.get_file(settings_path)
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8', errors='replace'))
            tab_size = settings.get('editor.tabSize', None)
    except Exception:
        pass
    return {'file_exists': file_exists, 'tab_size': tab_size}

def get_vscode_locale_and_settings__de00e991b11dacec5fa1917014fc1ed1_qw35sft2_973307ec(env, config: dict):
    """Read VS Code argv.json (locale) and settings.json (editor settings) from the VM."""
    result = {'locale': '', 'settings': {}}
    try:
        argv_bytes = env.controller.get_file('/home/user/.config/Code/argv.json')
        argv_data = json.loads(argv_bytes.decode('utf-8'))
        result['locale'] = argv_data.get('locale', '')
    except Exception as e:
        result['locale_error'] = str(e)
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings_data = json.loads(settings_bytes.decode('utf-8'))
        result['settings'] = settings_data
    except Exception:
        result['settings'] = {}
    return result

def get_workspace_and_fontsize__4cb0241f3e1c4176b65966dc96627af4_qw35sft2_73075398(env, config: dict):
    """Check workspace file existence and editor.fontSize in VS Code user settings."""
    ws_result = env.controller.run_bash_script('cat /home/user/project.code-workspace 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(ws_result, dict):
        ws_content = ws_result.get('output', '') or ws_result.get('stdout', '') or ''
    else:
        ws_content = str(ws_result)
    workspace_exists = '__NOT_FOUND__' not in ws_content and bool(ws_content.strip())
    settings_result = env.controller.run_bash_script('cat /home/user/.config/Code/User/settings.json 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(settings_result, dict):
        settings_content = settings_result.get('output', '') or settings_result.get('stdout', '') or ''
    else:
        settings_content = str(settings_result)
    font_size = None
    if '__NOT_FOUND__' not in settings_content and settings_content.strip():
        try:
            settings = json.loads(settings_content)
            font_size = settings.get('editor.fontSize')
        except (json.JSONDecodeError, ValueError):
            font_size = None
    return {'workspace_exists': workspace_exists, 'editor_font_size': font_size}

def get_vscode_py_indent_docstring__3d465870a880c6e96493a5637b57636e_qw35sft2_e21fc5b4(env, config: dict):
    """Read test.py from VM and check indentation of lines 2-10 and presence of a docstring."""
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/test.py')
        if not file_bytes:
            return {'error': 'File not found'}
        content = file_bytes.decode('utf-8', errors='replace')
        lines = content.splitlines()
        indent_ok = False
        if len(lines) >= 10:
            target_lines = lines[1:10]
            non_empty = [l for l in target_lines if l.strip()]
            indent_ok = bool(non_empty) and all((len(l) - len(l.lstrip()) >= 4 for l in non_empty))
        REQUIRED_DOCSTRING_TEXT = 'Sort the list in-place using bubble sort.'
        has_docstring = False
        in_func_body = False
        for l in lines:
            stripped = l.strip()
            if stripped.startswith('def ') and stripped.endswith(':'):
                in_func_body = True
                continue
            if in_func_body and ('"""' in stripped or "'''" in stripped):
                if REQUIRED_DOCSTRING_TEXT in stripped:
                    has_docstring = True
                break
            if in_func_body and stripped and (not stripped.startswith('#')) and ('"""' not in stripped):
                break
        return {'indent_ok': indent_ok, 'has_docstring': has_docstring}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_debug_focus_font__802820e3fd86caf9b92a9751803125a2_qw35sft2_77839c8f(env, config: dict):
    """Get debug.focusEditorOnBreak and editor.fontSize from settings.json."""
    try:
        file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings = json.loads(file_bytes.decode('utf-8'))
    except Exception:
        settings = {}
    return {'focus_editor_on_break': settings.get('debug.focusEditorOnBreak', True), 'font_size': settings.get('editor.fontSize', None)}

def get_vscode_settings_dict__8bfa6a8d7119af313a0660a3dea4ceab_qw35sft2_13b525eb(env, config: dict):
    """Read VS Code user settings.json and return parsed dict."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_keybindings__70885a60e67c4def20a9d0e3284c1793_qw35sft2_d6793f48(env, config: dict):
    """Read VS Code keybindings.json and return parsed entries."""
    keybindings_path = '/home/user/.config/Code/User/keybindings.json'
    file_bytes = env.controller.get_file(keybindings_path)
    if not file_bytes:
        return {'error': 'keybindings.json not found', 'entries': []}
    try:
        content = file_bytes.decode('utf-8')
        content = re.sub('//[^\\n]*', '', content)
        content = re.sub('/\\*.*?\\*/', '', content, flags=re.DOTALL)
        entries = json.loads(content.strip())
        return {'entries': entries if isinstance(entries, list) else []}
    except Exception as e:
        return {'error': str(e), 'entries': []}

def get_vscode_settings__fa93c32e8afd83b84bab5f585818b0fb_qw35sft2_e1ecdb7a(env, config: dict):
    """Read VS Code User settings.json from the VM and return its contents as a dict."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8')
        settings = json.loads(content)
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_ext_and_keybindings__44f31fc07357666acac95e9a419ee84f_qw35sft2_560113fa(env, config: dict):
    """Get extension list and keybindings.json content."""
    ext_list = ''
    try:
        response = requests.post(f'http://{env.vm_ip}:{env.server_port}/execute', json={'command': ['code', '--list-extensions'], 'shell': False}, timeout=30)
        if response.status_code == 200:
            ext_list = response.json().get('output', '')
    except Exception as e:
        logger_qw35sft2_d1773d.warning('Failed to get extension list: %s', e)
    keybindings = []
    try:
        kb_bytes = env.controller.get_file('/home/user/.config/Code/User/keybindings.json')
        if kb_bytes:
            raw = kb_bytes.decode('utf-8').strip()
            try:
                keybindings = json.loads(raw)
            except json.JSONDecodeError:
                lines = raw.splitlines()
                if lines and lines[0].startswith('//'):
                    keybindings = json.loads('\n'.join(lines[1:]))
    except Exception as e:
        logger_qw35sft2_d1773d.warning('Failed to read keybindings.json: %s', e)
    return {'ext_list': ext_list, 'keybindings': keybindings}

def get_vscode_theme_wordwrap__978e7cba9602912bc979b36fdc172116_qw35sft2_6665e744(env, config: dict):
    """Read VS Code settings.json and return colorTheme and editor.wordWrap."""
    path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'Settings file not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'colorTheme': settings.get('workbench.colorTheme'), 'wordWrap': settings.get('editor.wordWrap')}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_txt_content__004ad4b9f1b7c0d57d366b67195586e3_qw35sft2_8fb86838(env, config: dict):
    """Read vscode_replace_text.txt from the VM and return its text content."""
    path = config.get('path', '/home/user/Desktop/vscode_replace_text.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found or empty'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {str(e)}'}
    return {'content': content}

def get_workspace_folders__f72e9e4f450564b5101c9194862d205b_qw35sft2_32d45db6(env, config: dict):
    """Read VS Code workspace file and return normalized folder names."""
    workspace_path = config.get('path', '/home/user/project.code-workspace')
    file_bytes = env.controller.get_file(workspace_path)
    if not file_bytes:
        return {'error': 'Workspace file not found', 'folder_names': []}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        raw_paths = [f.get('path', '') for f in data.get('folders', [])]
        ws_dir = os.path.dirname(workspace_path)
        folder_names = []
        for p in raw_paths:
            if os.path.isabs(p):
                folder_names.append(os.path.basename(p.rstrip('/')))
            else:
                folder_names.append(os.path.basename(os.path.normpath(os.path.join(ws_dir, p))))
        return {'folder_names': folder_names, 'raw_paths': raw_paths}
    except Exception as e:
        return {'error': str(e), 'folder_names': []}

def get_vscode_settings__d45b70bac3694dd736206b10dff26100_qw35sft2_cf1fa267(env, config: dict):
    """Read VS Code user settings.json and return its parsed content."""
    path = config.get('path', '/home/user/.config/Code/User/settings.json')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8') if isinstance(file_bytes, bytes) else file_bytes
        return json.loads(content)
    except Exception as e:
        return {'error': str(e)}

def get_vscode_exclude_pytest__6a66711a30fb269a0e7d77ca8b497959_qw35sft2_2762ee26(env, config: dict):
    """Get VS Code files.exclude dict from settings.json."""
    files_exclude = {}
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            files_exclude = settings.get('files.exclude', {})
    except Exception as e:
        logger_qw35sft2_43fd17.warning('Failed to read settings.json: %s', e)
    return {'files_exclude': files_exclude}

def get_vscode_remove_and_add_keybinding__8bf5980383e9d80a451eac2b89ccf1df_qw35sft2_2666838f(env, config: dict):
    """Read keybindings.json: check ctrl+f removal for list.find AND ctrl+alt+e addition for Explorer."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/keybindings.json')
    if not file_bytes:
        return {'error': 'keybindings.json not found'}
    try:
        content = file_bytes.decode('utf-8').strip()
        if not content or content == '[]':
            return {'ctrlf_removed': False, 'explorer_shortcut_added': False}
        bindings = json.loads(content)
        ctrlf_removed = False
        explorer_shortcut_added = False
        for b in bindings:
            key = b.get('key', '').lower().replace(' ', '')
            cmd = b.get('command', '')
            when = b.get('when', '')
            if key == 'ctrl+f' and cmd == '-list.find' and ('listFocus' in when) and ('listSupportsFind' in when):
                ctrlf_removed = True
            if key == 'ctrl+alt+e' and cmd == 'workbench.view.explorer':
                explorer_shortcut_added = True
        return {'ctrlf_removed': ctrlf_removed, 'explorer_shortcut_added': explorer_shortcut_added}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_word_wrap_settings__35c11665b6e5c9d9c9abe7f855ccfe46_qw35sft2_80c5a18c(env, config: dict):
    """Read VS Code user settings.json and return wordWrapColumn and wordWrap values."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'wordWrapColumn': settings.get('editor.wordWrapColumn'), 'wordWrap': settings.get('editor.wordWrap')}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_locale_and_settings__e258fe9bb867a78f6ed6748a6bee2270_qw35sft2_d3d959ab(env, config: dict):
    """Read VS Code argv.json (locale) and settings.json (editor settings) from the VM."""
    result = {'locale': '', 'settings': {}}
    try:
        argv_bytes = env.controller.get_file('/home/user/.config/Code/argv.json')
        argv_data = json.loads(argv_bytes.decode('utf-8'))
        result['locale'] = argv_data.get('locale', '')
    except Exception as e:
        result['locale_error'] = str(e)
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings_data = json.loads(settings_bytes.decode('utf-8'))
        result['settings'] = settings_data
    except Exception:
        result['settings'] = {}
    return result

def get_vscode_ext_and_settings__e2317900e4eb6fe8b0284d167aa38270_qw35sft2_2f9de2bd(env, config: dict):
    """Get VS Code extension list and user settings.json."""
    import json
    ext_result = env.controller.run_bash_script('code --list-extensions 2>/dev/null', timeout=30)
    if isinstance(ext_result, dict):
        ext_output = ext_result.get('output', '') or ext_result.get('stdout', '') or ''
    else:
        ext_output = str(ext_result) if ext_result else ''
    settings_result = env.controller.run_bash_script('cat /home/user/.config/Code/User/settings.json 2>/dev/null || echo "{}"', timeout=30)
    if isinstance(settings_result, dict):
        settings_str = settings_result.get('output', '') or settings_result.get('stdout', '') or '{}'
    else:
        settings_str = str(settings_result) if settings_result else '{}'
    try:
        settings = json.loads(settings_str.strip())
    except Exception:
        settings = {}
    return {'extensions': ext_output.strip().lower(), 'settings': settings}

def get_workspace_and_extension__bead7b293cda299d38413ae3ff873cc8_qw35sft2_9fcd0186(env, config: dict):
    """Check workspace file existence and whether Python extension is installed."""
    ws_result = env.controller.run_bash_script('cat /home/user/project.code-workspace 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(ws_result, dict):
        ws_content = ws_result.get('output', '') or ws_result.get('stdout', '') or ''
    else:
        ws_content = str(ws_result)
    workspace_exists = '__NOT_FOUND__' not in ws_content and bool(ws_content.strip())
    ext_result = env.controller.run_bash_script('code --list-extensions 2>/dev/null || echo "__ERROR__"', timeout=30)
    if isinstance(ext_result, dict):
        ext_content = ext_result.get('output', '') or ext_result.get('stdout', '') or ''
    else:
        ext_content = str(ext_result)
    extensions = []
    if '__ERROR__' not in ext_content:
        extensions = [line.strip().lower() for line in ext_content.splitlines() if line.strip()]
    return {'workspace_exists': workspace_exists, 'installed_extensions': extensions}

def get_vscode_py_function_file__9e0f9c32752df4ad0956b05cabebbc98_qw35sft2_e6aaeebb(env, config: dict):
    """Get test.py from Desktop and extract function definition info."""
    import re
    file_path = config.get('path', '/home/user/Desktop/test.py')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'content': '', 'has_function': False, 'function_name': None}
        content = file_bytes.decode('utf-8', errors='replace')
        match = re.search('^\\s*def\\s+(\\w+)\\s*\\(', content, re.MULTILINE)
        has_function = match is not None
        function_name = match.group(1) if match else None
        return {'exists': True, 'content': content, 'has_function': has_function, 'function_name': function_name}
    except Exception as e:
        return {'exists': False, 'content': '', 'has_function': False, 'function_name': None, 'error': str(e)}

def get_vscode_workspace_in_storage__5752e6d2cd4f13c8d24f076655028c23_qw35sft2_a6c9eb8a(env, config: dict):
    """Check if VS Code has opened a specific workspace by inspecting workspaceStorage."""
    target = config.get('target_workspace', 'workspace1.code-workspace')
    script = f'find /home/user/.config/Code/User/workspaceStorage/ -name "workspace.json" 2>/dev/null | xargs grep -l "{target}" 2>/dev/null | head -1'
    result = env.controller.run_bash_script(script, timeout=15)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    found = bool(output and output.strip())
    return {'workspace_opened': found, 'target': target}

def get_vscode_py_indent_blankline__2c836d0eca9976b79042c49c51fe71cf_qw35sft2_67fb627b(env, config: dict):
    """Read test.py from VM and check indentation of lines 2-10 and blank line after function def."""
    try:
        file_bytes = env.controller.get_file('/home/user/Desktop/test.py')
        if not file_bytes:
            return {'error': 'File not found'}
        content = file_bytes.decode('utf-8', errors='replace')
        lines = content.splitlines()
        n_line = next((l for l in lines if l.lstrip() == 'n = len(alist)'), None)
        for_j_line = next((l for l in lines if l.lstrip().startswith('for j in range')), None)
        indent_ok = False
        if n_line is not None and for_j_line is not None:
            n_indent = len(n_line) - len(n_line.lstrip())
            for_j_indent = len(for_j_line) - len(for_j_line.lstrip())
            indent_ok = n_indent >= 4 and for_j_indent >= 4
        has_blank_after_def = False
        for idx, line in enumerate(lines):
            stripped = line.strip()
            if stripped.startswith('def ') and stripped.endswith(':'):
                if idx + 1 < len(lines) and lines[idx + 1].strip() == '':
                    has_blank_after_def = True
                break
        return {'indent_ok': indent_ok, 'has_blank_after_def': has_blank_after_def}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_settings_dict__9ec2182c98e7d1d45bbbb3ae2d0db42e_qw35sft2_be433eca(env, config: dict):
    """Read VS Code user settings.json and return parsed dict."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_debug_focus_wrap__09cc1b24020e8f7c0f7b48c6a4d96413_qw35sft2_4aeffc35(env, config: dict):
    """Get debug.focusEditorOnBreak and editor.wordWrap from settings.json."""
    try:
        file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings = json.loads(file_bytes.decode('utf-8'))
    except Exception:
        settings = {}
    return {'focus_editor_on_break': settings.get('debug.focusEditorOnBreak', True), 'word_wrap': settings.get('editor.wordWrap', 'off')}

def get_vscode_settings__a4b8eefcd32d2ac2b259c28c71814550_qw35sft2_af664dff(env, config: dict):
    """Read VS Code User settings.json from the VM and return its contents as a dict."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8')
        settings = json.loads(content)
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_workspace_folders__d5491cdef8176b887dce317a57139150_qw35sft2_81e65fe6(env, config: dict):
    """Read VS Code workspace file and return normalized folder names."""
    workspace_path = config.get('path', '/home/user/project.code-workspace')
    file_bytes = env.controller.get_file(workspace_path)
    if not file_bytes:
        return {'error': 'Workspace file not found', 'folder_names': []}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        raw_paths = [f.get('path', '') for f in data.get('folders', [])]
        ws_dir = os.path.dirname(workspace_path)
        folder_names = []
        for p in raw_paths:
            if os.path.isabs(p):
                folder_names.append(os.path.basename(p.rstrip('/')))
            else:
                folder_names.append(os.path.basename(os.path.normpath(os.path.join(ws_dir, p))))
        return {'folder_names': folder_names, 'raw_paths': raw_paths}
    except Exception as e:
        return {'error': str(e), 'folder_names': []}

def get_vscode_settings__5cd158cc213aefb25242b372a3862ff3_qw35sft2_3d23d8ed(env, config: dict):
    """Read VS Code settings.json and return parsed settings dict."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'settings.json not found', 'settings': {}}
    try:
        content = file_bytes.decode('utf-8')
        content = re.sub('//[^\\n]*', '', content)
        content = re.sub('/\\*.*?\\*/', '', content, flags=re.DOTALL)
        settings = json.loads(content.strip())
        return {'settings': settings if isinstance(settings, dict) else {}}
    except Exception as e:
        return {'error': str(e), 'settings': {}}

def get_vscode_txt_content__076bfd2d3ed2cd7e42a937a1dd80c8a9_qw35sft2_b2456901(env, config: dict):
    """Read vscode_replace_text.txt from the VM and return its text content."""
    path = config.get('path', '/home/user/Desktop/vscode_replace_text.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found or empty'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {str(e)}'}
    return {'content': content}

def get_vscode_wrap_and_format__95a613ea8825af09f40057c978f02ad7_qw35sft2_f6281378(env, config: dict):
    """Read VS Code user settings.json and return wordWrapColumn and formatOnSave values."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'wordWrapColumn': settings.get('editor.wordWrapColumn'), 'formatOnSave': settings.get('editor.formatOnSave')}
    except Exception as e:
        return {'error': str(e)}

def get_workspace_and_wordwrap__1713b5c17b096bf267a520ed66d388f0_qw35sft2_f0f8a503(env, config: dict):
    """Check workspace file existence and editor.wordWrap in VS Code user settings."""
    ws_result = env.controller.run_bash_script('cat /home/user/project.code-workspace 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(ws_result, dict):
        ws_content = ws_result.get('output', '') or ws_result.get('stdout', '') or ''
    else:
        ws_content = str(ws_result)
    workspace_exists = '__NOT_FOUND__' not in ws_content and bool(ws_content.strip())
    settings_result = env.controller.run_bash_script('cat /home/user/.config/Code/User/settings.json 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(settings_result, dict):
        settings_content = settings_result.get('output', '') or settings_result.get('stdout', '') or ''
    else:
        settings_content = str(settings_result)
    word_wrap = None
    if '__NOT_FOUND__' not in settings_content and settings_content.strip():
        try:
            settings = json.loads(settings_content)
            word_wrap = settings.get('editor.wordWrap')
        except (json.JSONDecodeError, ValueError):
            word_wrap = None
    return {'workspace_exists': workspace_exists, 'editor_word_wrap': word_wrap}

def get_vscode_exclude_autosave__7797789b3c2086d41b9d5f1699195599_qw35sft2_4b32b4aa(env, config: dict):
    """Get VS Code files.exclude and files.autoSave from settings.json."""
    files_exclude = {}
    auto_save = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            files_exclude = settings.get('files.exclude', {})
            auto_save = settings.get('files.autoSave')
    except Exception as e:
        logger_qw35sft2_c6cea1.warning('Failed to read settings.json: %s', e)
    return {'files_exclude': files_exclude, 'auto_save': auto_save}

def get_vscode_desktop_file_comment__eea98fddd6e0176c917193edbee2b203_qw35sft2_3e7e2657(env, config: dict):
    """Get test.py from Desktop and return its first non-empty line."""
    file_path = config.get('path', '/home/user/Desktop/test.py')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'first_line': '', 'content': ''}
        content = file_bytes.decode('utf-8', errors='replace')
        lines = [l for l in content.splitlines() if l.strip()]
        first_line = lines[0].strip() if lines else ''
        return {'exists': True, 'first_line': first_line, 'content': content}
    except Exception as e:
        return {'exists': False, 'first_line': '', 'content': '', 'error': str(e)}

def get_vscode_ext_and_settings__88ac2a64a7ee030ff707715de339ea34_qw35sft2_c88a3e48(env, config: dict):
    """Get VS Code extension list and user settings.json."""
    import json
    ext_result = env.controller.run_bash_script('code --list-extensions 2>/dev/null', timeout=30)
    if isinstance(ext_result, dict):
        ext_output = ext_result.get('output', '') or ext_result.get('stdout', '') or ''
    else:
        ext_output = str(ext_result) if ext_result else ''
    settings_result = env.controller.run_bash_script('cat /home/user/.config/Code/User/settings.json 2>/dev/null || echo "{}"', timeout=30)
    if isinstance(settings_result, dict):
        settings_str = settings_result.get('output', '') or settings_result.get('stdout', '') or '{}'
    else:
        settings_str = str(settings_result) if settings_result else '{}'
    try:
        settings = json.loads(settings_str.strip())
    except Exception:
        settings = {}
    return {'extensions': ext_output.strip().lower(), 'settings': settings}

def get_vscode_settings_dict__d122d575a22853e47d56ee92b6c18378_qw35sft2_b8e62d25(env, config: dict):
    """Read VS Code user settings.json and return parsed dict."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_txt_content__e470774862e9eb18ed4fb0ad458ab39f_qw35sft2_abefc668(env, config: dict):
    """Read vscode_replace_text.txt from the VM and return its text content."""
    path = config.get('path', '/home/user/Desktop/vscode_replace_text.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found or empty'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {str(e)}'}
    return {'content': content}

def get_vscode_open_workspace__c37846152daf06e05c18277c5b1f636d_qw35sft2_5b8c70c2(env, config: dict):
    """Detect the workspace folder currently open in VS Code via window title."""
    try:
        cmd_result = env.controller.run_bash_script('xdotool search --onlyvisible --class \'Code\' 2>/dev/null | head -10 | while read wid; do xdotool getwindowname "$wid" 2>/dev/null; done', timeout=15)
        output = ''
        if isinstance(cmd_result, dict):
            output = (cmd_result.get('output') or cmd_result.get('stdout') or '').strip()
        elif isinstance(cmd_result, (bytes, str)):
            output = (cmd_result if isinstance(cmd_result, str) else cmd_result.decode('utf-8', errors='ignore')).strip()
        for line in output.splitlines():
            line = line.strip()
            if ' - Visual Studio Code' in line:
                workspace_name = line.split(' - Visual Studio Code')[0].strip()
                if workspace_name.startswith('● '):
                    workspace_name = workspace_name[2:].strip()
                if ' - ' in workspace_name:
                    workspace_name = workspace_name.split(' - ')[-1].strip()
                return {'workspace_name': workspace_name.lower()}
        tree = env.controller.get_accessibility_tree()
        if tree:
            if 'PROJECT' in tree:
                return {'workspace_name': 'project'}
        return {'workspace_name': None}
    except Exception as e:
        return {'workspace_name': None, 'error': str(e)}

def get_workspace_folders__edd7ecbec0169642eef724f574a58652_qw35sft2_04296562(env, config: dict):
    """Read VS Code workspace file and return normalized folder names."""
    workspace_path = config.get('path', '/home/user/project.code-workspace')
    file_bytes = env.controller.get_file(workspace_path)
    if not file_bytes:
        return {'error': 'Workspace file not found', 'folder_names': []}
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        raw_paths = [f.get('path', '') for f in data.get('folders', [])]
        ws_dir = os.path.dirname(workspace_path)
        folder_names = []
        for p in raw_paths:
            if os.path.isabs(p):
                folder_names.append(os.path.basename(p.rstrip('/')))
            else:
                folder_names.append(os.path.basename(os.path.normpath(os.path.join(ws_dir, p))))
        return {'folder_names': folder_names, 'raw_paths': raw_paths}
    except Exception as e:
        return {'error': str(e), 'folder_names': []}

def get_vscode_workspace_state__f91feeeb0a889e0b6b274960d3d06c7d_qw35sft2_cc26b79e(env, config: dict):
    """Check if VS Code workspace file was saved at /home/user/project.code-workspace."""
    result = env.controller.run_bash_script('cat /home/user/project.code-workspace 2>/dev/null || echo "__NOT_FOUND__"', timeout=10)
    if isinstance(result, dict):
        content = result.get('output', '') or result.get('stdout', '') or ''
    else:
        content = str(result)
    if '__NOT_FOUND__' in content or not content.strip():
        return {'workspace_exists': False, 'valid_json': False}
    try:
        data = json.loads(content)
        return {'workspace_exists': True, 'valid_json': True, 'data': data}
    except (json.JSONDecodeError, ValueError):
        return {'workspace_exists': True, 'valid_json': False, 'data': {}}

def get_vscode_tab_size__015bdcf4a8f96215b5579fc53db2d526_qw35sft2_860be523(env, config: dict):
    """Read VS Code user settings.json and return the tabSize value."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    try:
        file_bytes = env.controller.get_file(settings_path)
        if not file_bytes:
            return {'error': 'Settings file not found'}
        settings = json.loads(file_bytes.decode('utf-8'))
        return {'tabSize': settings.get('editor.tabSize')}
    except Exception as e:
        return {'error': str(e)}

def get_vscode_debug_focus_terminal__29baecbd3036f18eff818341b1c0e4e8_qw35sft2_0d7e8ce8(env, config: dict):
    """Get debug.focusEditorOnBreak and terminal.integrated.cursorBlinking from settings.json."""
    try:
        file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        settings = json.loads(file_bytes.decode('utf-8'))
    except Exception:
        settings = {}
    return {'focus_editor_on_break': settings.get('debug.focusEditorOnBreak', True), 'cursor_blinking': settings.get('terminal.integrated.cursorBlinking', False)}

def get_vscode_exclude_fontsize__33cb4b35f94afe16fa1fe71880064b0b_qw35sft2_745131bd(env, config: dict):
    """Get VS Code files.exclude and editor.fontSize from settings.json."""
    files_exclude = {}
    font_size = None
    try:
        settings_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
        if settings_bytes:
            settings = json.loads(settings_bytes.decode('utf-8'))
            files_exclude = settings.get('files.exclude', {})
            font_size = settings.get('editor.fontSize')
    except Exception as e:
        logger_qw35sft2_65fcf1.warning('Failed to read settings.json: %s', e)
    return {'files_exclude': files_exclude, 'font_size': font_size}

def get_vscode_settings__d97bffb72b4b8c27def4b087597b01a0_qw35sft2_5da34fe4(env, config: dict):
    """Read VS Code User settings.json from the VM and return its contents as a dict."""
    settings_path = '/home/user/.config/Code/User/settings.json'
    file_bytes = env.controller.get_file(settings_path)
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        content = file_bytes.decode('utf-8')
        settings = json.loads(content)
        return settings
    except Exception as e:
        return {'error': str(e)}

def get_vscode_ext_list__521eb992bf97de2d79e4aa6c081db684_qw35sft2_67841176(env, config: dict):
    """Get installed VS Code extension list."""
    result = env.controller.run_bash_script('code --list-extensions 2>/dev/null', timeout=30)
    if isinstance(result, dict):
        output = result.get('output', '') or result.get('stdout', '') or ''
    else:
        output = str(result) if result else ''
    return {'extensions': output.strip().lower()}

def get_vscode_txt_content__d87d9b6cb6d6a7fc9d5d36ba0c96e0de_qw35sft2_c1aa3798(env, config: dict):
    """Read vscode_replace_text.txt from the VM and return its text content."""
    path = config.get('path', '/home/user/Desktop/vscode_replace_text.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found or empty'}
    try:
        content = file_bytes.decode('utf-8')
    except Exception as e:
        return {'error': f'Decode error: {str(e)}'}
    return {'content': content}

def get_vscode_active_file__6a582969bc6d0b9624887969279510e3_qw35sft2_2cbb6469(env, config: dict):
    """Detect the currently active editor file in VS Code via window title."""
    try:
        cmd_result = env.controller.run_bash_script('xdotool search --onlyvisible --class \'Code\' 2>/dev/null | head -10 | while read wid; do xdotool getwindowname "$wid" 2>/dev/null; done', timeout=15)
        output = ''
        if isinstance(cmd_result, dict):
            output = (cmd_result.get('output') or cmd_result.get('stdout') or '').strip()
        elif isinstance(cmd_result, (bytes, str)):
            output = (cmd_result if isinstance(cmd_result, str) else cmd_result.decode('utf-8', errors='ignore')).strip()
        for line in output.splitlines():
            line = line.strip()
            if ' - Visual Studio Code' in line:
                title_part = line.split(' - Visual Studio Code')[0].strip()
                if title_part.startswith('● '):
                    title_part = title_part[2:].strip()
                if ' - ' in title_part:
                    active_file = title_part.split(' - ')[0].strip()
                else:
                    active_file = None
                return {'active_file': active_file}
        tree = env.controller.get_accessibility_tree()
        if tree:
            target = config.get('expected_file', '')
            if target and target in tree:
                return {'active_file': target}
        return {'active_file': None}
    except Exception as e:
        return {'active_file': None, 'error': str(e)}

def get_vscode_settings_dict__7e25498cca269e458737cc32ebf8beb0_qw35sft2_477d89cd(env, config: dict):
    """Read VS Code user settings.json and return parsed dict."""
    file_bytes = env.controller.get_file('/home/user/.config/Code/User/settings.json')
    if not file_bytes:
        return {'error': 'settings.json not found'}
    try:
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        return {'error': str(e)}
