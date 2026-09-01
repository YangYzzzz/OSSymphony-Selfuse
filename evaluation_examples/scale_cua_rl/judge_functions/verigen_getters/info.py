"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_compose_window_title__df6702fb4aa059fb4a86e567f6f9f62d', 'get_screen_lock_settings__39d04f914d2dbd8627ddcd7f791e785b', 'get_settings_screen_dims__29161befcd93aa27778463661523fc23', 'get_settings_window_size__e47348d3e93e8a085dab2f1753ad88d5', 'get_settings_window_size__a8bc55c18275bf61543bfdfa6bb3d447_qw35sft2_ee274776', 'get_screen_blank_delay__40046a4f8d5906bbddcb5abf1fee6b07_qw35sft2_2b1a48b1', 'get_screen_lock_state__74e051f75a4f9998d114f6407d55ddd5_qw35sft2_ead426cc', 'get_accessibility_large_text_screen_keyboard__0c6c70afb2a250c0e6dcd4a5ca3ca3b8_qw35sft2_aa9d77a1', 'get_screen_lock_full_state__8b9cedf68783a02fb2242a963784d1c2_qw35sft2_9f505d12', 'get_accessibility_large_text_screen_reader__e5c088a08a33da472617cdfc723aca05_qw35sft2_7f704d7d', 'get_screen_blank_delay__3696e962dd7f0719e2ddde3289beabec_qw35sft2_21db7fcf', 'get_notif_dnd_lockscreen__7a1296a554f735dccd7fa86e30d68dbd_qw35sft2_af587089', 'get_screen_lock_full_state__b5802fb719515ed2e2b308492efceb55_qw35sft2_c07db47b', 'get_notif_screenlock__0ee0cafb23f5579534516156a3b19db3_qw35sft2_1bd901f0', 'get_screen_lock_full_state__5e9c7d265f840de83da743ff562a1457_qw35sft2_4efffc9b', 'get_restore_wallpaper__354747c46b7fabc52ac1f7b0a9e657e8_qw35sft2_ff8213ec', 'get_desktop_wallpaper_uri__d24bdc521260cf2eb34948f145dea565_qw35sft2_fbd5df14']

def get_compose_window_title__df6702fb4aa059fb4a86e567f6f9f62d(env, config: dict):
    """Get the Thunderbird compose window title using xdotool."""
    try:
        result = env.controller.run_bash_script("xdotool search --name 'Write:' getwindowname 2>/dev/null | head -1", timeout=30)
        output = ''
        if isinstance(result, dict):
            output = result.get('output', result.get('stdout', ''))
        elif isinstance(result, str):
            output = result
        return {'window_title': output.strip()}
    except Exception as e:
        return {'error': str(e)}

def get_screen_lock_settings__39d04f914d2dbd8627ddcd7f791e785b(env, config: dict):
    """Get screen lock enabled status and idle delay settings."""
    try:
        idle_result = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=30)
        lock_result = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=30)
        idle_output = idle_result.get('output', '').strip() if isinstance(idle_result, dict) else str(idle_result).strip()
        lock_output = lock_result.get('output', '').strip() if isinstance(lock_result, dict) else str(lock_result).strip()
        return {'idle_delay': idle_output, 'lock_enabled': lock_output}
    except Exception as e:
        return {'error': str(e)}

def get_settings_screen_dims__29161befcd93aa27778463661523fc23(env, config: dict):
    """Get SCREEN_WIDTH and SCREEN_HEIGHT from settings.py."""
    path = config.get('path', '/home/user/Desktop/tetris/settings.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    width_match = re.search('SCREEN_WIDTH\\s*=\\s*(\\d+)', content)
    height_match = re.search('SCREEN_HEIGHT\\s*=\\s*(\\d+)', content)
    return {'screen_width': int(width_match.group(1)) if width_match else None, 'screen_height': int(height_match.group(1)) if height_match else None}

def get_settings_window_size__e47348d3e93e8a085dab2f1753ad88d5(env, config: dict):
    """Get WIDTH and HEIGHT values from settings.py."""
    try:
        result = env.controller.run_bash_script('python3 -c "exec(open(\'/home/user/Desktop/snake/settings.py\').read()); print(WIDTH); print(HEIGHT)"', timeout=30)
        output = result.get('output', '').strip()
        lines = output.split('\n')
        if len(lines) >= 2:
            width = int(float(lines[0].strip()))
            height = int(float(lines[1].strip()))
            return {'width': width, 'height': height}
        return {'error': 'Could not parse WIDTH/HEIGHT from output', 'raw': output}
    except Exception as e:
        return {'error': str(e)}

def get_settings_window_size__a8bc55c18275bf61543bfdfa6bb3d447_qw35sft2_ee274776(env, config: dict):
    """Read settings.py and extract WIDTH and HEIGHT values."""
    import re
    file_bytes = env.controller.get_file('/home/user/Desktop/snake/settings.py')
    if not file_bytes:
        return {'error': 'settings.py not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    match = re.search('^WIDTH\\s*,\\s*HEIGHT\\s*=\\s*(\\d+)\\s*,\\s*(\\d+)', content, re.MULTILINE)
    if match:
        return {'width': int(match.group(1)), 'height': int(match.group(2))}
    w_match = re.search('^WIDTH\\s*=\\s*(\\d+)', content, re.MULTILINE)
    h_match = re.search('^HEIGHT\\s*=\\s*(\\d+)', content, re.MULTILINE)
    if w_match and h_match:
        return {'width': int(w_match.group(1)), 'height': int(h_match.group(1))}
    return {'error': 'WIDTH/HEIGHT not found in settings.py'}

def get_screen_blank_delay__40046a4f8d5906bbddcb5abf1fee6b07_qw35sft2_2b1a48b1(env, config: dict):
    """Get the screen blank (idle-delay) setting value in seconds."""
    r = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    if isinstance(r, dict):
        output = r.get('output', r.get('stdout', '')).strip()
    else:
        output = str(r).strip()
    parts = output.split()
    delay = int(parts[-1]) if parts else -1
    return {'idle_delay': delay}

def get_screen_lock_state__74e051f75a4f9998d114f6407d55ddd5_qw35sft2_ead426cc(env, config: dict):
    """Get automatic screen lock enabled state via gsettings."""
    result = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=10)
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', '')).strip()
    else:
        output = str(result).strip()
    return {'lock_enabled': output}

def get_accessibility_large_text_screen_keyboard__0c6c70afb2a250c0e6dcd4a5ca3ca3b8_qw35sft2_aa9d77a1(env, config: dict):
    """Get large-text and screen-keyboard-enabled states via gsettings."""
    vm_ip = env.vm_ip
    port = env.server_port
    result = {}
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.interface', 'large-text'], 'shell': False})
    if resp1.status_code == 200:
        result['large_text'] = resp1.json()['output'].strip() == 'true'
    else:
        result['large_text'] = None
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.applications', 'screen-keyboard-enabled'], 'shell': False})
    if resp2.status_code == 200:
        result['screen_keyboard'] = resp2.json()['output'].strip() == 'true'
    else:
        result['screen_keyboard'] = None
    return result

def get_screen_lock_full_state__8b9cedf68783a02fb2242a963784d1c2_qw35sft2_9f505d12(env, config: dict):
    """Get screen lock enabled state and blank screen delay via gsettings."""
    lock_raw = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=10)
    delay_raw = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    lock_str = lock_raw.get('output', lock_raw.get('stdout', '')).strip() if isinstance(lock_raw, dict) else str(lock_raw).strip()
    delay_str = delay_raw.get('output', delay_raw.get('stdout', '')).strip() if isinstance(delay_raw, dict) else str(delay_raw).strip()
    return {'lock_enabled': lock_str, 'idle_delay': delay_str}

def get_accessibility_large_text_screen_reader__e5c088a08a33da472617cdfc723aca05_qw35sft2_7f704d7d(env, config: dict):
    """Get large-text and screen-reader-enabled states via gsettings."""
    vm_ip = env.vm_ip
    port = env.server_port
    result = {}
    resp1 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.interface', 'large-text'], 'shell': False})
    if resp1.status_code == 200:
        result['large_text'] = resp1.json()['output'].strip() == 'true'
    else:
        result['large_text'] = None
    resp2 = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': ['gsettings', 'get', 'org.gnome.desktop.a11y.applications', 'screen-reader-enabled'], 'shell': False})
    if resp2.status_code == 200:
        result['screen_reader'] = resp2.json()['output'].strip() == 'true'
    else:
        result['screen_reader'] = None
    return result

def get_screen_blank_delay__3696e962dd7f0719e2ddde3289beabec_qw35sft2_21db7fcf(env, config: dict):
    """Get the screen blank (idle-delay) setting value in seconds."""
    r = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    if isinstance(r, dict):
        output = r.get('output', r.get('stdout', '')).strip()
    else:
        output = str(r).strip()
    parts = output.split()
    delay = int(parts[-1]) if parts else -1
    return {'idle_delay': delay}

def get_notif_dnd_lockscreen__7a1296a554f735dccd7fa86e30d68dbd_qw35sft2_af587089(env, config: dict):
    """Get DND (show-banners) and lock-screen notification settings."""
    r1 = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-banners', timeout=15)
    if isinstance(r1, dict):
        show_banners = r1.get('output', r1.get('stdout', '')).strip()
    else:
        show_banners = str(r1).strip()
    r2 = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-in-lock-screen', timeout=15)
    if isinstance(r2, dict):
        show_in_lock_screen = r2.get('output', r2.get('stdout', '')).strip()
    else:
        show_in_lock_screen = str(r2).strip()
    return {'show_banners': show_banners, 'show_in_lock_screen': show_in_lock_screen}

def get_screen_lock_full_state__b5802fb719515ed2e2b308492efceb55_qw35sft2_c07db47b(env, config: dict):
    """Get screen lock enabled state and blank screen delay via gsettings."""
    lock_raw = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=10)
    delay_raw = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    lock_enabled = lock_raw.get('output', lock_raw.get('stdout', '')).strip() if isinstance(lock_raw, dict) else str(lock_raw).strip()
    idle_delay = delay_raw.get('output', delay_raw.get('stdout', '')).strip() if isinstance(delay_raw, dict) else str(delay_raw).strip()
    return {'lock_enabled': lock_enabled, 'idle_delay': idle_delay}

def get_notif_screenlock__0ee0cafb23f5579534516156a3b19db3_qw35sft2_1bd901f0(env, config: dict):
    """Get DND (show-banners) and screen lock enabled settings."""
    r1 = env.controller.run_bash_script('gsettings get org.gnome.desktop.notifications show-banners', timeout=15)
    show_banners = r1.get('output', r1.get('stdout', '')).strip() if isinstance(r1, dict) else str(r1).strip()
    r2 = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=15)
    lock_enabled = r2.get('output', r2.get('stdout', '')).strip() if isinstance(r2, dict) else str(r2).strip()
    return {'show_banners': show_banners, 'lock_enabled': lock_enabled}

def get_screen_lock_full_state__5e9c7d265f840de83da743ff562a1457_qw35sft2_4efffc9b(env, config: dict):
    """Get screen lock enabled state and blank screen delay via gsettings."""
    raw_lock = env.controller.run_bash_script('gsettings get org.gnome.desktop.screensaver lock-enabled', timeout=10)
    lock_enabled = raw_lock.get('output', raw_lock.get('stdout', '')).strip() if isinstance(raw_lock, dict) else str(raw_lock).strip()
    raw_delay = env.controller.run_bash_script('gsettings get org.gnome.desktop.session idle-delay', timeout=10)
    idle_delay = raw_delay.get('output', raw_delay.get('stdout', '')).strip() if isinstance(raw_delay, dict) else str(raw_delay).strip()
    return {'lock_enabled': lock_enabled, 'idle_delay': idle_delay}

def get_restore_wallpaper__354747c46b7fabc52ac1f7b0a9e657e8_qw35sft2_ff8213ec(env, config: dict):
    """Check if poster is on Desktop and set as GNOME desktop wallpaper."""
    vm_ip = env.vm_ip
    port = env.server_port
    try:
        ls_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'ls /home/user/Desktop/ 2>/dev/null', 'shell': True}, timeout=15)
        if ls_resp.status_code != 200:
            return {'error': f'ls HTTP {ls_resp.status_code}'}
        ls_output = ls_resp.json().get('output', '') or ''
        desktop_files = [f.strip() for f in ls_output.strip().split('\n') if f.strip()]
        wp_resp = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': 'gsettings get org.gnome.desktop.background picture-uri 2>/dev/null', 'shell': True}, timeout=15)
        wallpaper_uri = ''
        if wp_resp.status_code == 200:
            wallpaper_uri = (wp_resp.json().get('output', '') or '').strip().strip("'")
        return {'file_on_desktop': 'poster_party_night.webp' in desktop_files, 'wallpaper_uri': wallpaper_uri}
    except Exception as e:
        logger_qw35sft2_2c81bb.error('get_restore_wallpaper__354747c46b7fabc52ac1f7b0a9e657e8 error: %s', e)
        return {'error': str(e)}

def get_desktop_wallpaper_uri__d24bdc521260cf2eb34948f145dea565_qw35sft2_fbd5df14(env, config: dict):
    """Get the current desktop wallpaper URI via gsettings."""
    cmd = 'gsettings get org.gnome.desktop.background picture-uri'
    result = env.controller.run_bash_script(cmd, timeout=15)
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    uri = output.strip().strip('\'"')
    return {'wallpaper_uri': uri}
