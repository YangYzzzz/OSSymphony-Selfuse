"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_gimp_multi_config__06220b8da042a0d79908bcbdd2480275', 'get_image_dimensions__58760a5b814fae5bb8129a050acf3bc8', 'get_image_properties__6fac91d7fdc9ac567c2522f5862ea0c3', 'get_image_files_in_dir__c59a92777846cb030643927ac9628267', 'get_image_dimensions__7d65af646c8700a29016dc9b48791bd4', 'get_image_flip_data__30cadb0ec7eca28df693cf924ab88301', 'get_gif_file_info__53c5a4efb8d4cc457e973819b6da6e9c', 'get_image_properties__f8142f5206c8e9d081140bc00d3fc065', 'get_image_color_mode__3e8ba2ce3f71bc35828deb0d6a741a53', 'get_image_dimensions__d91421f1c316cdb714dbaed12454e46f', 'get_image_dimensions__ebe9791a027ef956f5ed9c5e287e6589', 'get_image_properties__4cd9609153e6f966f651b58aa2377fc7', 'get_image_info__3254fa401a56ff7bd854b02b7cc3c8d5', 'get_gimp_sharpness_data__dd81583a369d4d4452910cdac2a642dc', 'get_image_dimensions__96e0bf39947996b933028cd995042bff', 'get_gimp_theme_and_icons__48364d2e461d23c72442c1de0ff1aa61', 'get_jpeg_file_check__eecb20058330e9d3e08bcd27beed34f8', 'get_image_color_mode__742fdbd7c8ad904b152a719424c640d8', 'get_image_mode_check__a066a98104bf8bbb65a3648dc0f66596', 'get_image_file_exists__9f8821d93c1286af4ab340177f286f49', 'get_image_dimensions__7532d635f451c203dc4c9c0128dc01d3', 'get_image_dimensions__73568344a1f5a1b5f117e9d5ad98afd8', 'get_gimp_export_state__c940821317f4393b9e9b23183cb5bab0', 'get_gimp_png_dimensions__a3a90a0c275eafc9a08c6844b32106ee_qw35sft2_5ef93f8a', 'get_gimp_undo_and_theme__402f86b3aa5767626ba0a47b48426a4d_qw35sft2_37149e62', 'get_gimp_png_dimensions__6daba05d23c98408363b21827c57395a_qw35sft2_24bd5ba1', 'get_gimp_image_mode__e4220acd868ea9e3c374fad54e986f17_qw35sft2_3dcba9a4', 'get_image_mode__9561abdc46f323aa49f39811909461f6_qw35sft2_e1c0b3d5', 'get_gimp_img_dims__ee2a931afc08caeb2f7cbb79edbe900c_qw35sft2_5bb80250', 'get_gimp_xcf_layers__f6dc05e0a2477162d5d2aaf28dcb399e_qw35sft2_6e8436d4', 'get_gimp_image_props__fc1506ee66b785bd74f37812296aa386_qw35sft2_7cfbd2a7', 'get_image_size_stat__e32eea47feb3c826ca65513ea1a8c80e_qw35sft2_73d96f59', 'get_gimp_darktable_plugin__93f45a58587e292e529b56db6ae27226_qw35sft2_d9994b22', 'get_gimp_saturation_contrast__918f3cc47704c578c4667b1ca6019596_qw35sft2_d82ed289', 'get_gimp_config_multi__d66c76f0f1e5e42f832534b4f9451e5c_qw35sft2_c82d0c91', 'get_gimp_dims_brightness__0ffd5d6c48e00cc11530a6062d573f43_qw35sft2_90d64ff9', 'get_gimp_image_mode__a22a6a74169a7c2b7bd4f15c461d97a5_qw35sft2_adb429cd', 'get_gimp_image_thumbnail__4ddc2b8f762091345213742b2e539174_qw35sft2_33ee6ccc', 'get_gimp_session_single_window__940ba8eb8583aa4e66192a1528358acd_qw35sft2_9b24c412', 'get_gimp_image_mode__b63a2518d07e9d19e5e5bc7d70fabd81_qw35sft2_019efa2b', 'get_gimp_saturation_stats__b674a0dc75a50659a17535e66cb88a78_qw35sft2_a56e2f2c', 'get_gimp_yellow_centroid__466f916fdda48299f33abc0cff1e195a_qw35sft2_60e5fb31', 'get_gimp_image_dims__ffabe808aeb093a1184f15812c26092c_qw35sft2_e25cffe9', 'get_gimp_image_and_layer__5e787196d91537456fde84b3a66eeecb_qw35sft2_6a04fddf', 'get_gimp_png_bg_pixel__4d0454ba1a7fe5978c1747f21d797308_qw35sft2_074a772f', 'get_gimp_hd_png__6d65b36cd0a438b39e8f21308e345a4a_qw35sft2_848c1efa', 'get_gimp_image_size__08c1ae085b68de96a31c7d77e1455141_qw35sft2_87a5475c', 'get_gimp_xcf_layers__ad9b8a9ed8fc4caff86597f320aba5f2_qw35sft2_a7561ace', 'get_gimp_xcf_layers__f3ec988a99b0f642c7ac2a56df923954_qw35sft2_59992397', 'get_gimp_rot180_check__359fadf8cec20093badb6005d3119d0c_qw35sft2_b61f5722', 'get_gimp_process__be67663f10a797bb181cd6fa722efa34_qw35sft2_7486e82e', 'get_gimp_saturation_size__e189188d82a47045e4921db4c0cbb19b_qw35sft2_1d6c1177', 'get_gimp_grayscale_brightness__72b5517a5a800033bb98cd0fbc3da308_qw35sft2_b713d12c', 'get_gimp_image_props__e5dd5bd57fb559e638682724445a4fb8_qw35sft2_de8a7ed3', 'get_gimp_file_exists__6e337deb47fb4ded97926f2aaf4149ad_qw35sft2_4ca91268', 'get_image_brightness_stat__1dd8caa614c12c2ce7c1b58be9a02c11_qw35sft2_dce61cdc', 'get_gimp_jpeg_stats__99df0d1fa420bfe8f4d58bfec01f5388_qw35sft2_9f161e6b', 'get_image_dimensions__1649a3c87969acc3b9b0ad261438802f_qw35sft2_8ed03cf1', 'get_gimp_image_dims__2254c4ac73b2e0e2d0bc95e32c93006c_qw35sft2_ae629eeb', 'get_gimp_grayscale_check__c8310c6f68cacbb7f806947f886f9156_qw35sft2_75ba4ee6', 'get_gimp_image_mode__5cf8f74046cc654cd234f12e0bd56fa2_qw35sft2_e1f7e6c9', 'get_gimp_icon_theme__8b7e0a81347ead101665acba492eaaa3_qw35sft2_75b7fca1', 'get_gimp_png_orientation__65c7fe1c7d5d90a17513cb0cbcf64984_qw35sft2_38622937', 'get_gimp_layer_bbox__6fac91d7fdc9ac567c2522f5862ea0c3_qw35sft2_ce514f78', 'get_gimp_exported_png__2034cd2081886ce4d15c41eaddab9687_qw35sft2_569c72d7', 'get_gimp_png_bg_pixel__c69df4a705c1504ceebda9205ba51ae8_qw35sft2_123f6c98', 'get_gimp_hflip_check__e47abffb5b090802d85e5308a4df8da8_qw35sft2_689caf68', 'get_gimp_theme_setting__a2417c567c82dd8e5e291c848fcf86af_qw35sft2_31a64438', 'get_gimp_saturation_stats__6b90b70b3d1b161e089c6f4f8582392c_qw35sft2_ad606690', 'get_gimp_canvas_and_layer__0d4c5e593127b30f4279483bc536e453_qw35sft2_6c86c613', 'get_desktop_jpeg_exists__e78c01185325a38b78f35c525682fa2c_qw35sft2_33aee3e6', 'get_gimp_image_color_info__4728404e320da0722e3f0bbf27779384_qw35sft2_7a1c2b28', 'get_gimp_jpeg_export__3494183243a5e3cf64667682e9d5cf69_qw35sft2_b423750f', 'get_gimp_gray_stats__a80023f11b098859420cd69b0067e271_qw35sft2_e03c66fa', 'get_desktop_xcf__30619c16638f3c14c2b868710d3b511a_qw35sft2_3341f65e', 'get_gimp_config_multi__59d11816792f7732bb8631499d5a3dca_qw35sft2_434bfb14', 'get_gimp_flipped_file__f5b72d9abaa33261d0da153e643d90eb_qw35sft2_02e67ab3', 'get_gimp_image_flip__917488a366fb0a5d7f8700214812c742_qw35sft2_ea5a8b0d', 'get_gimp_rotated_dims__480c196a41fcbd8c06328eee76532eca_qw35sft2_a1188a1a', 'get_gimp_file_exists__9dede0efb5d54dc2e4df137a06f2c4ba_qw35sft2_e1eea9d6', 'get_gimp_dark_theme__48287b5138535fb1248d29dc8ceabce5_qw35sft2_c404d8a1', 'get_image_dimensions__9d59aa3517eac66c5b7e40985af8d11e_qw35sft2_82cd2c39', 'get_gimp_png_text_pos__f282f693e16d842a8fd79d730244ff72_qw35sft2_f5a11855', 'get_image_brightness_stat__53b38a939b32dab621a5954331033a3c_qw35sft2_3c0a82cc', 'get_gimp_yellow_centroid__569dedcfdc3e7d291a20f0e6aca641c0_qw35sft2_8d4fd89e', 'get_gimp_image_dims__2bd82e314cabe351963c3af4824b9b87_qw35sft2_c504650a', 'get_gimp_layer_bbox_opacity__c6fafc0c96511b4f7b94b1c390c2f8bd_qw35sft2_a1553309', 'get_gimp_png_bg_pixel__3cdf6e4d2fe41e5dc87ba95dd55c83c0_qw35sft2_b80b5fe7', 'get_gimp_vflip_check__e3f424dc6a1401f7936bf6fab235a0b4_qw35sft2_f99490de', 'get_gimp_icon_theme__a072fad36e1102073ae062e8b644b246_qw35sft2_6918c39c', 'get_gimp_image_props__94a488f0af5ba53f580224747ebe24fe_qw35sft2_3884b707', 'get_gimp_xcf_layers__5ffee09eda4febddd3875d5ca422ec9b_qw35sft2_0ad9a860', 'get_gimp_dims_brightness__1abf0259addada4e9320a412de93a720_qw35sft2_f0440aa5', 'get_gimp_config_multi__d46e54f2b8fa1eafde00e6af910958bf_qw35sft2_3aa5062d', 'get_gimp_theme_setting__ff19a7440acd5edb92781b332298a214_qw35sft2_f3f4120e', 'get_gimp_scale_stats__4ce9d4657814d8ed93df9fa8b61db4cf_qw35sft2_8a4fa70a', 'get_gimp_png_export__fcf5b1c8c41f94f2ef23ac5bc8b1ca54_qw35sft2_ebcdfd42', 'get_gimp_saturation_sharpness__cc63708ed8a0e82770a04258c8e27328_qw35sft2_c9a74c7a', 'get_gimp_docks_window__e872a1deec6ced899c5abf40a9c68092_qw35sft2_42f2e145', 'get_gimp_png_square__ba207a15040c6bc16b5ad1659309ab17_qw35sft2_f3c5241b', 'get_gimp_yellow_centroid__bc6d20fd4fedacbbf14f4629287f2cc1_qw35sft2_56a18ff2', 'get_gimp_file_size__5c243e9c006808209adf482ad7ac1fe7_qw35sft2_cadec8c0', 'get_jpeg_file_exists__132eea53d5785c870b331091d6b9fa18_qw35sft2_b2f768ea', 'get_gimp_image_dims__392ddc5ce54a4061c2326df12808cf9c_qw35sft2_21889611', 'get_gimp_layer_bbox__4cd9609153e6f966f651b58aa2377fc7_qw35sft2_58886e85', 'get_gimp_png_bg_pixel__06570640015ea4861f9ec59a895e97a7_qw35sft2_175b87d7', 'get_three_images_brightness__b3992f4af3f3949be116d5f6a9289f79_qw35sft2_c55a41c0', 'get_gimp_jpeg_info__d6fa459b37f87da875d5b6370f3cc58a_qw35sft2_b8f15e4c', 'get_gimp_image_mode__f54b84b51bb414e1baf7bdf3d2e0e400_qw35sft2_37cb28b8', 'get_gimp_xcf_layers__413739fbbe273da47f148920d82ae49b_qw35sft2_3e65ef23', 'get_gimp_image_props__857686b0194013b75de5a057cd91e5ad_qw35sft2_8e186571', 'get_desktop_image_list__1dbc4bff650a401c3959c27fb5e70015_qw35sft2_e03db9b1', 'get_transition_and_png__dd35705ee09fd488827a6afdd7cbab81_qw35sft2_03943646', 'get_docx_image_and_alignment__b349befcb19ab6c9d751d2833a5ceb8c_qw35sft2_8ed4cb3f', 'get_docx_image_count__ed4f0e4477e03c56fe74b28d9d6a3444_qw35sft2_acb79876', 'get_docx_image_and_fontsize__285ee55215732bb31d8255226d83a755_qw35sft2_abebb093', 'get_docx_image_and_pdf__cba3b2d7a3f4e1f51b9a84ac71e8e772_qw35sft2_e2f152b3', 'get_docx_image_and_pagenumbers__6390ff8ee787412726544ea6ba43db0d_qw35sft2_fd1dbc8e', 'get_small_jpeg_size__98edc79a225cc1dd6b5d8a23724241ab_qw35sft2_1b132276', 'get_image_dimensions__b392913b04d99d3e7513b78f2dedf15d_qw35sft2_078d6e64', 'get_image_size_state__8cc8ca1163a74ceeec7ab0fdb85713f3_qw35sft2_5465ba2b', 'get_image_dimensions__89f1952738bca9658039f047c420ba2d_qw35sft2_38f42e78', 'get_jpeg_file_size__19b267768d7fcf0b434b3cca9c02b15d_qw35sft2_a7b7be93', 'get_jpg_move_state__7427978e92f6fc0e5652a4713261a5a8_qw35sft2_821f0d4a', 'get_jpg_png_copy_state__ec3ddc36152b3d2fb4aee4f74f969e31_qw35sft2_2c06a286', 'get_jpg_copy_with_count__58fb65f54a1152c1f3aecb552e15d932_qw35sft2_28376ab0', 'get_jpg_filelist__b885c2e3b122c9010091a38454018836_qw35sft2_048d9885', 'get_vacation_jpg_copy__91c6f86e45abe96db716a4e3d1072be2_qw35sft2_fe78b515', 'get_vlc_image_adjust__1e9b265786c948c671445fe3130d0b36_qw35sft2_687fc6f9', 'get_vlc_image_adjust__33b1c5d9e144110cf196db624efc6d81_qw35sft2_573798c0', 'get_vlc_image_adjust__493da9b8d1a20a26a98a5cc60b751e43_qw35sft2_2106ab38', 'get_vlc_image_adjust__c451d2a7e7a5b8e510cce0c34da4325d_qw35sft2_3cf3258a', 'get_vlc_image_adjust__ab50e89f0aef704ab57d4c9d2579417f_qw35sft2_83e6738b']

def get_gimp_multi_config__06220b8da042a0d79908bcbdd2480275(env, config: dict):
    """Get multiple key-value pairs from GIMP gimprc config file."""
    import tempfile
    import os
    try:
        file_bytes = env.controller.get_file(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))
    except Exception:
        try:
            file_bytes = env.controller.get_file(os.path.expanduser('~/.gimp-2.10/gimprc'))
        except Exception:
            return {'error': 'Could not find gimprc'}
    if not file_bytes:
        return {'error': 'gimprc file is empty or not found'}
    content = file_bytes.decode('utf-8', errors='replace')
    keys_to_check = config.get('keys', [])
    result = {}
    for key in keys_to_check:
        pattern = f'\\({re.escape(key)}\\s+(.+?)\\)'
        match = re.search(pattern, content)
        if match:
            result[key] = match.group(1).strip().strip('"')
        else:
            result[key] = None
    return result

def get_image_dimensions__58760a5b814fae5bb8129a050acf3bc8(env, config: dict):
    """Download image file from VM and return its dimensions."""
    from PIL import Image
    path = config.get('path', '/home/user/Desktop/character.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'path': path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'format': img.format}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_image_properties__6fac91d7fdc9ac567c2522f5862ea0c3(env, config: dict):
    """Get image properties including dimensions."""
    from PIL import Image
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (full_width, full_height) = img.size
        return {'width': full_width, 'height': full_height}
    finally:
        os.unlink(tmp_path)

def get_image_files_in_dir__c59a92777846cb030643927ac9628267(env, config: dict):
    """List image files in a directory on the VM."""
    target_dir = config.get('path', '/home/user/Pictures')
    result = env.controller.run_bash_script(f'find "{target_dir}" -maxdepth 1 -type f \\( -iname "*.png" -o -iname "*.jpg" -o -iname "*.jpeg" -o -iname "*.gif" -o -iname "*.bmp" \\) 2>/dev/null | sort', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    if not output:
        return {'image_files': [], 'count': 0}
    files = [f for f in output.split('\n') if f.strip()]
    return {'image_files': files, 'count': len(files)}

def get_image_dimensions__7d65af646c8700a29016dc9b48791bd4(env, config: dict):
    """Get image dimensions from a file on the VM."""
    from PIL import Image
    import io
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    try:
        img = Image.open(io.BytesIO(file_bytes))
        (width, height) = img.size
        return {'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_image_flip_data__30cadb0ec7eca28df693cf924ab88301(env, config: dict):
    """
    Getter that fetches two images from VM:
    - result_path: the output image (pic.jpg)
    - original_path: the original image
    Returns both images as raw bytes for comparison in metric.
    """
    result_path = config.get('result_path', '/home/user/Desktop/pic.jpg')
    original_path = config.get('original_path', '')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        return {'error': f'Result file not found: {result_path}'}
    original_bytes = env.controller.get_file(original_path)
    if not original_bytes:
        return {'error': f'Original file not found: {original_path}'}
    try:
        from PIL import Image
        from skimage.metrics import structural_similarity as ssim
        import numpy as np
        result_img = Image.open(io.BytesIO(result_bytes)).convert('RGB')
        original_img = Image.open(io.BytesIO(original_bytes)).convert('RGB')
        flipped_original = original_img.transpose(Image.FLIP_LEFT_RIGHT)
        if result_img.size != flipped_original.size:
            result_img = result_img.resize(flipped_original.size, Image.LANCZOS)
        result_arr = np.array(result_img)
        flipped_arr = np.array(flipped_original)
        similarity = ssim(result_arr, flipped_arr, channel_axis=2)
        return {'ssim': float(similarity), 'result_size': result_img.size, 'original_size': original_img.size}
    except Exception as e:
        return {'error': f'Image comparison failed: {str(e)}'}

def get_gif_file_info__53c5a4efb8d4cc457e973819b6da6e9c(env, config: dict):
    """Get info about a GIF file on the VM to verify it was created correctly."""
    file_path = config.get('path', '/home/user/Desktop/src_output.gif')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    is_gif = file_bytes[:3] == b'GIF'
    gif_version = file_bytes[:6].decode('ascii', errors='replace') if is_gif else ''
    file_size = len(file_bytes)
    width = 0
    height = 0
    if is_gif and len(file_bytes) >= 10:
        width = int.from_bytes(file_bytes[6:8], byteorder='little')
        height = int.from_bytes(file_bytes[8:10], byteorder='little')
    return {'exists': True, 'is_gif': is_gif, 'gif_version': gif_version, 'file_size': file_size, 'width': width, 'height': height}

def get_image_properties__f8142f5206c8e9d081140bc00d3fc065(env, config: dict):
    """Get layer dimensions from the layer_dims.txt file written by Script-Fu batch."""
    file_path = config.get('path', '/home/user/Desktop/layer_dims.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found: ' + file_path}
    content = file_bytes.decode('utf-8', errors='replace').strip()
    if not content:
        return {'error': 'Empty dimensions file'}
    layers = {}
    for line in content.split('\n'):
        line = line.strip()
        if not line:
            continue
        parts = line.rsplit(' ', 2)
        if len(parts) == 3:
            try:
                name = parts[0]
                w = int(parts[1])
                h = int(parts[2])
                layers[name] = {'width': w, 'height': h}
            except ValueError:
                continue
    if not layers:
        return {'error': 'Could not parse layer dimensions', 'raw': content}
    target_layer = config.get('target_layer', 'dog')
    for (name, dims) in layers.items():
        if target_layer.lower() in name.lower():
            return {'width': dims['width'], 'height': dims['height'], 'layer_name': name}
    for (name, dims) in layers.items():
        if 'background' not in name.lower():
            return {'width': dims['width'], 'height': dims['height'], 'layer_name': name}
    return {'error': 'Target layer not found', 'layers': list(layers.keys())}

def get_image_color_mode__3e8ba2ce3f71bc35828deb0d6a741a53(env, config: dict):
    """Get the color mode of an image file from the VM."""
    from PIL import Image
    import numpy as np
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        mode = img.mode
        is_grayscale = False
        if mode in ('L', 'LA', 'P'):
            is_grayscale = True
        elif mode in ('RGB', 'RGBA'):
            arr = np.array(img)
            if arr.shape[2] >= 3:
                is_grayscale = np.allclose(arr[:, :, 0], arr[:, :, 1], atol=2) and np.allclose(arr[:, :, 1], arr[:, :, 2], atol=2)
        return {'mode': mode, 'is_grayscale': is_grayscale, 'width': img.size[0], 'height': img.size[1]}
    finally:
        os.unlink(tmp_path)

def get_image_dimensions__d91421f1c316cdb714dbaed12454e46f(env, config: dict):
    """Get image dimensions from a file on the VM."""
    from PIL import Image
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (w, h) = img.size
        return {'width': w, 'height': h}
    finally:
        os.unlink(tmp_path)

def get_image_dimensions__ebe9791a027ef956f5ed9c5e287e6589(env, config: dict):
    """Get image dimensions from a file on the VM."""
    from PIL import Image
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (w, h) = img.size
        return {'width': w, 'height': h}
    finally:
        os.unlink(tmp_path)

def get_image_properties__4cd9609153e6f966f651b58aa2377fc7(env, config: dict):
    """Get layer dimensions from a GIMP XCF file using Script-Fu batch mode.

    Reads the saved XCF and queries layer dimensions directly, avoiding
    the need to export (which flattens layers and loses transparency info).
    """
    xcf_path = config['path']
    layer_name = config.get('layer_name', 'dog')
    dims_file = '/tmp/gimp_layer_dims.txt'
    script_file = '/tmp/gimp_query_dims.scm'
    scm_script = '(let* ((image (car (gimp-file-load RUN-NONINTERACTIVE "{xcf}" "{xcf}")))\n       (num (car (gimp-image-get-layers image)))\n       (arr (cadr (gimp-image-get-layers image)))\n       (port (open-output-file "{out}")))\n  (let loop ((i 0))\n    (if (< i num)\n      (let* ((layer (vector-ref arr i))\n             (name (car (gimp-item-get-name layer)))\n             (w (car (gimp-drawable-width layer)))\n             (h (car (gimp-drawable-height layer))))\n        (display name port)\n        (display ":" port)\n        (display w port)\n        (display "x" port)\n        (display h port)\n        (newline port)\n        (loop (+ i 1)))))\n  (close-output-port port)\n  (gimp-image-delete image))\n'.format(xcf=xcf_path, out=dims_file)
    env.controller.run_bash_script(f"cat > {script_file} << 'GIMP_SCM_EOF'\n{scm_script}GIMP_SCM_EOF", timeout=10)
    env.controller.run_bash_script(f'rm -f {dims_file}', timeout=5)
    env.controller.run_bash_script(f'''gimp -i -b '(load "{script_file}")' -b '(gimp-quit 0)' 2>/dev/null''', timeout=60)
    file_bytes = env.controller.get_file(dims_file)
    if not file_bytes:
        return {'error': 'Failed to get layer dimensions'}
    content = file_bytes.decode('utf-8').strip()
    layers = {}
    for line in content.split('\n'):
        line = line.strip()
        if ':' in line and 'x' in line:
            (name_part, dims_part) = line.rsplit(':', 1)
            parts = dims_part.split('x', 1)
            if len(parts) == 2:
                try:
                    layers[name_part.strip()] = {'width': int(parts[0]), 'height': int(parts[1])}
                except ValueError:
                    continue
    target = layer_name.lower()
    for (name, dims) in layers.items():
        if name.lower() == target or target in name.lower():
            return {'width': dims['width'], 'height': dims['height'], 'layer_name': name, 'all_layers': layers}
    return {'error': f'Layer "{layer_name}" not found', 'all_layers': layers}

def get_image_info__3254fa401a56ff7bd854b02b7cc3c8d5(env, config: dict):
    """Get image file info (existence, dimensions, dominant color) from VM."""
    path = config.get('path', '')
    try:
        result = env.controller.run_bash_script(f'identify -format "%wx%h" "{path}" 2>/dev/null && echo "" && stat -c "%s" "{path}" 2>/dev/null && echo "" && convert "{path}" -resize 1x1 txt:- 2>/dev/null', timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        if not output:
            return {'error': 'File not found or not an image', 'exists': False}
        lines = output.strip().split('\n')
        if len(lines) >= 2:
            dims = lines[0].strip()
            size_bytes = lines[1].strip()
            parts = dims.split('x')
            if len(parts) == 2:
                info = {'exists': True, 'width': int(parts[0]), 'height': int(parts[1]), 'file_size': int(size_bytes), 'is_white': False}
                if len(lines) >= 3:
                    color_line = lines[-1].strip()
                    if '#FFFFFF' in color_line.upper() or '#FEFEFE' in color_line.upper():
                        info['is_white'] = True
                    elif '(255,255,255)' in color_line or '(65535,65535,65535)' in color_line:
                        info['is_white'] = True
                    else:
                        import re
                        rgb_match = re.search('\\((\\d+),(\\d+),(\\d+)\\)', color_line)
                        if rgb_match:
                            (r, g, b) = (int(rgb_match.group(1)), int(rgb_match.group(2)), int(rgb_match.group(3)))
                            if r > 60000 and g > 60000 and (b > 60000):
                                info['is_white'] = True
                            elif r >= 250 and g >= 250 and (b >= 250) and (r <= 255):
                                info['is_white'] = True
                return info
        return {'error': 'Could not parse image info', 'exists': False}
    except Exception as e:
        return {'error': str(e), 'exists': False}

def get_gimp_sharpness_data__dd81583a369d4d4452910cdac2a642dc(env, config: dict):
    """
    Fetches both edited and original images from the VM, computes sharpness
    (Laplacian variance) and structural similarity for comparison.

    config keys:
      - edited_path: path to the edited (sharpened) image on VM
      - original_path: path to the original image on VM
    """
    try:
        from PIL import Image
        from skimage.metrics import structural_similarity as ssim
        import cv2
    except ImportError as e:
        return {'error': f'Missing dependency: {e}'}
    edited_path = config.get('edited_path', '')
    original_path = config.get('original_path', '')
    edited_bytes = env.controller.get_file(edited_path)
    if not edited_bytes:
        return {'error': f'Edited file not found: {edited_path}'}
    original_bytes = env.controller.get_file(original_path)
    if not original_bytes:
        return {'error': f'Original file not found: {original_path}'}
    edited_tmp = None
    original_tmp = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(edited_bytes)
            edited_tmp = tmp.name
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(original_bytes)
            original_tmp = tmp.name
        edited_img = cv2.imread(edited_tmp)
        original_img = cv2.imread(original_tmp)
        if edited_img is None:
            return {'error': 'Failed to decode edited image'}
        if original_img is None:
            return {'error': 'Failed to decode original image'}
        edited_gray = cv2.cvtColor(edited_img, cv2.COLOR_BGR2GRAY)
        original_gray = cv2.cvtColor(original_img, cv2.COLOR_BGR2GRAY)
        edited_laplacian_var = cv2.Laplacian(edited_gray, cv2.CV_64F).var()
        original_laplacian_var = cv2.Laplacian(original_gray, cv2.CV_64F).var()
        if edited_gray.shape != original_gray.shape:
            original_gray_resized = cv2.resize(original_gray, (edited_gray.shape[1], edited_gray.shape[0]))
        else:
            original_gray_resized = original_gray
        sim_score = ssim(edited_gray, original_gray_resized)
        return {'edited_sharpness': float(edited_laplacian_var), 'original_sharpness': float(original_laplacian_var), 'structural_similarity': float(sim_score)}
    finally:
        if edited_tmp and os.path.exists(edited_tmp):
            os.unlink(edited_tmp)
        if original_tmp and os.path.exists(original_tmp):
            os.unlink(original_tmp)

def get_image_dimensions__96e0bf39947996b933028cd995042bff(env, config: dict):
    """Download image file from VM and return its dimensions."""
    from PIL import Image
    path = config.get('path', '/home/user/Desktop/character_icon.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'path': path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'format': img.format}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_theme_and_icons__48364d2e461d23c72442c1de0ff1aa61(env, config: dict):
    """Get both theme and icon-theme from GIMP config file."""
    import tempfile
    import os
    import re
    try:
        file_bytes = env.controller.get_file('/home/user/.config/GIMP/2.10/gimprc')
        if not file_bytes:
            file_bytes = env.controller.get_file('/home/user/.gimp-2.10/gimprc')
        if not file_bytes:
            return {'error': 'gimprc not found'}
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False, mode='wb') as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, 'r') as f:
                content = f.read()
        finally:
            os.unlink(tmp_path)
        result = {'theme': None, 'icon_theme': None}
        theme_match = re.search('\\(theme\\s+"([^"]+)"\\)', content)
        if theme_match:
            result['theme'] = theme_match.group(1)
        icon_match = re.search('\\(icon-theme\\s+"([^"]+)"\\)', content)
        if icon_match:
            result['icon_theme'] = icon_match.group(1)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_jpeg_file_check__eecb20058330e9d3e08bcd27beed34f8(env, config: dict):
    """Check if a JPEG file exists on the VM and is valid."""
    from PIL import Image
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'is_jpeg': False, 'width': 0, 'height': 0}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_jpeg = img.format == 'JPEG'
        (w, h) = img.size
        return {'exists': True, 'is_jpeg': is_jpeg, 'width': w, 'height': h}
    except Exception:
        return {'exists': True, 'is_jpeg': False, 'width': 0, 'height': 0}
    finally:
        os.unlink(tmp_path)

def get_image_color_mode__742fdbd7c8ad904b152a719424c640d8(env, config: dict):
    """Get the color mode of an image file from the VM."""
    from PIL import Image
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.size[0], 'height': img.size[1]}
    finally:
        os.unlink(tmp_path)

def get_image_mode_check__a066a98104bf8bbb65a3648dc0f66596(env, config: dict):
    """Check the color mode of an image file on the VM."""
    import tempfile, os
    from PIL import Image, ImageChops
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    suffix = '.' + config['path'].rsplit('.', 1)[-1] if '.' in config['path'] else '.jpeg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        mode = img.mode
        (width, height) = img.size
        is_grayscale = mode in ('L', 'LA')
        if not is_grayscale and mode == 'RGB':
            (r, g, b) = img.split()
            diff_rg = ImageChops.difference(r, g)
            diff_rb = ImageChops.difference(r, b)
            is_grayscale = diff_rg.getbbox() is None and diff_rb.getbbox() is None
        return {'mode': mode, 'is_grayscale': bool(is_grayscale), 'width': width, 'height': height}
    finally:
        os.unlink(tmp_path)

def get_image_file_exists__9f8821d93c1286af4ab340177f286f49(env, config: dict):
    """Check if an image file exists and get its size."""
    path = config.get('path', '')
    result = env.controller.run_bash_script(f"test -f '{path}' && stat --format='%s' '{path}' 2>/dev/null || echo 'NOT_FOUND'", timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'file_size': 0}
    try:
        size = int(output)
        return {'exists': True, 'file_size': size}
    except ValueError:
        return {'exists': False, 'file_size': 0}

def get_image_dimensions__7532d635f451c203dc4c9c0128dc01d3(env, config: dict):
    """Get image dimensions from a file on the VM."""
    from PIL import Image
    import io
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    try:
        img = Image.open(io.BytesIO(file_bytes))
        (width, height) = img.size
        return {'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_image_dimensions__73568344a1f5a1b5f117e9d5ad98afd8(env, config: dict):
    """Get image dimensions from a file on the VM."""
    from PIL import Image
    import io
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    try:
        img = Image.open(io.BytesIO(file_bytes))
        (width, height) = img.size
        return {'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_gimp_export_state__c940821317f4393b9e9b23183cb5bab0(env, config: dict):
    """
    Get the state of a GIMP-exported image: check if file exists and get dimensions.
    Returns a dict with 'exists' (bool) and 'width'/'height' (int or None).
    """
    file_path = config.get('path', '')
    check_result = env.controller.run_bash_script(f'test -f "{file_path}" && echo "EXISTS" || echo "MISSING"', timeout=10)
    file_exists = 'EXISTS' in check_result.get('output', '')
    if not file_exists:
        return {'exists': False, 'width': None, 'height': None}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        (width, height) = img.size
        return {'exists': True, 'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        return {'exists': True, 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_dimensions__a3a90a0c275eafc9a08c6844b32106ee_qw35sft2_5ef93f8a(env, config: dict):
    """Download an image file from the VM and return its dimensions."""
    import io
    path = config.get('path', '/home/user/Desktop/canvas.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        return {'width': img.width, 'height': img.height, 'format': img.format}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_gimp_undo_and_theme__402f86b3aa5767626ba0a47b48426a4d_qw35sft2_37149e62(env, config: dict):
    """Read undo-levels and theme from GIMP gimprc config file."""
    gimp_config_paths = ['/root/.config/GIMP/2.10/gimprc', '/home/user/.config/GIMP/2.10/gimprc']
    content = None
    for path in gimp_config_paths:
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            content = file_bytes.decode('utf-8', errors='ignore')
            break
    if content is None:
        return {'error': 'gimprc not found', 'undo_levels': None, 'theme': None}
    m_undo = re.search('\\(undo-levels\\s+([^\\)]+)\\)', content)
    undo_levels = m_undo.group(1).strip().strip('"') if m_undo else None
    m_theme = re.search('\\(theme\\s+([^\\)]+)\\)', content)
    theme = m_theme.group(1).strip().strip('"') if m_theme else None
    return {'undo_levels': undo_levels, 'theme': theme}

def get_gimp_png_dimensions__6daba05d23c98408363b21827c57395a_qw35sft2_24bd5ba1(env, config: dict):
    """Download the exported PNG and return its pixel dimensions."""
    import tempfile, os
    vm_path = config.get('path', '/home/user/Desktop/orange_background.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'File not found: ' + vm_path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        width, height = img.size
        return {'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_mode__e4220acd868ea9e3c374fad54e986f17_qw35sft2_3dcba9a4(env, config: dict):
    """Read export.jpg from the VM and return its PIL image mode."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/export.jpg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None}
    suffix = os.path.splitext(file_path)[-1] or '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'mode': None}
    finally:
        os.unlink(tmp_path)

def get_image_mode__9561abdc46f323aa49f39811909461f6_qw35sft2_e1c0b3d5(env, config: dict):
    """Get the PIL mode (e.g. 'L', 'RGB', 'RGBA') of the image on the Desktop."""
    path = config.get('path', '/home/user/Desktop/dog_with_background.png')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'File not found', 'mode': None}
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(file_bytes))
        return {'mode': img.mode}
    except Exception as e:
        return {'error': str(e), 'mode': None}

def get_gimp_img_dims__ee2a931afc08caeb2f7cbb79edbe900c_qw35sft2_5bb80250(env, config: dict):
    """
    Download an image file from the VM and return its pixel dimensions.
    Used to verify scaling/resize operations that produce specific canvas sizes.
    """
    import io
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL/Pillow not available'}
    path = config.get('path', '/home/user/Desktop/berry.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': f'File not found: {path}'}
    try:
        img = Image.open(io.BytesIO(file_bytes))
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e)}

def get_gimp_xcf_layers__f6dc05e0a2477162d5d2aaf28dcb399e_qw35sft2_6e8436d4(env, config: dict):
    """Parse XCF file from VM and return layer names and count."""
    import struct
    xcf_path = config.get('path', '/home/user/Desktop/white_background.xcf')
    try:
        file_bytes = env.controller.get_file(xcf_path)
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}
    if not file_bytes:
        return {'error': 'File not found or empty', 'layer_names': [], 'layer_count': 0}
    data = bytes(file_bytes)
    if not data.startswith(b'gimp xcf '):
        return {'error': 'Not a valid XCF file', 'layer_names': [], 'layer_count': 0}
    try:
        null_pos = data.index(b'\x00', 9)
        version_str = data[9:null_pos].decode('ascii', errors='replace')
        pos = null_pos + 1
        if pos + 12 > len(data):
            return {'error': 'File too short', 'layer_names': [], 'layer_count': 0}
        pos += 12
        while pos + 8 <= len(data):
            prop_type = struct.unpack('>I', data[pos:pos + 4])[0]
            prop_size = struct.unpack('>I', data[pos + 4:pos + 8])[0]
            pos += 8 + prop_size
            if prop_type == 0:
                break
        try:
            ver_num = int(version_str.lstrip('v') or 0)
        except Exception:
            ver_num = 0
        offset_size = 8 if ver_num >= 11 else 4
        offset_fmt = '>Q' if ver_num >= 11 else '>I'
        layer_offsets = []
        while pos + offset_size <= len(data):
            offset = struct.unpack(offset_fmt, data[pos:pos + offset_size])[0]
            pos += offset_size
            if offset == 0:
                break
            layer_offsets.append(offset)
        layer_names = []
        for off in layer_offsets:
            if off + 16 > len(data):
                continue
            name_off = off + 12
            name_len = struct.unpack('>I', data[name_off:name_off + 4])[0]
            name_off += 4
            if name_len == 0 or name_off + name_len > len(data):
                continue
            name_bytes = data[name_off:name_off + name_len]
            name = name_bytes.rstrip(b'\x00').decode('utf-8', errors='replace')
            layer_names.append(name)
        return {'layer_names': layer_names, 'layer_count': len(layer_names)}
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}

def get_gimp_image_props__fc1506ee66b785bd74f37812296aa386_qw35sft2_7cfbd2a7(env, config: dict):
    """Get image mode and dimensions from exported PNG file on the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/rotated_palette_computer.png')
    try:
        from PIL import Image
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'file not found or empty'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            img = Image.open(tmp_path)
            result = {'mode': img.mode, 'width': img.size[0], 'height': img.size[1]}
            img.close()
        finally:
            os.unlink(tmp_path)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_image_size_stat__e32eea47feb3c826ca65513ea1a8c80e_qw35sft2_73d96f59(env, config: dict):
    """Get width and height of an image file from the VM."""
    import tempfile, os
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL not available'}
    vm_path = config.get('path', '')
    if not vm_path:
        return {'error': 'No path specified in config'}
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': f'File not found or empty: {vm_path}', 'file_found': False}
    suffix = os.path.splitext(vm_path)[1] or '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        width, height = img.size
        return {'width': width, 'height': height, 'file_found': True}
    except Exception as e:
        return {'error': str(e), 'file_found': False}
    finally:
        os.unlink(tmp_path)

def get_gimp_darktable_plugin__93f45a58587e292e529b56db6ae27226_qw35sft2_d9994b22(env, config: dict):
    """Check if the darktable GIMP plugin (file-darktable) is installed."""
    plugin_path = '/usr/lib/gimp/2.0/plug-ins/file-darktable'
    result = env.controller.run_bash_script(f"test -f '{plugin_path}' && echo 'exists' || echo 'missing'", timeout=10)
    if isinstance(result, dict):
        stdout = result.get('output', result.get('stdout', '')).strip()
    else:
        stdout = str(result).strip()
    plugin_present = stdout == 'exists'
    pkg_result = env.controller.run_bash_script("dpkg -l darktable 2>/dev/null | grep -c '^ii' || echo '0'", timeout=15)
    if isinstance(pkg_result, dict):
        pkg_stdout = pkg_result.get('output', pkg_result.get('stdout', '')).strip()
    else:
        pkg_stdout = str(pkg_result).strip()
    pkg_installed = pkg_stdout.strip() == '1'
    return {'plugin_present': plugin_present, 'pkg_installed': pkg_installed, 'plugin_path': plugin_path}

def get_gimp_saturation_contrast__918f3cc47704c578c4667b1ca6019596_qw35sft2_d82ed289(env, config: dict):
    """Download edited image from VM and return mean HSV saturation and grayscale contrast std."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    file_bytes = env.controller.get_file('/home/user/Desktop/woman_sitting_by_the_tree2.png')
    if not file_bytes:
        return {'error': 'File not found on VM'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_rgb = img.convert('RGB')
        arr = np.array(img_rgb).astype(float) / 255.0
        r, g, b = (arr[:, :, 0], arr[:, :, 1], arr[:, :, 2])
        maxc = np.maximum(np.maximum(r, g), b)
        minc = np.minimum(np.minimum(r, g), b)
        sat = np.where(maxc > 0, (maxc - minc) / maxc, 0.0)
        img_gray = img.convert('L')
        gray_arr = np.array(img_gray).astype(float)
        contrast_std = float(gray_arr.std())
        return {'saturation_mean': float(sat.mean()), 'contrast_std': contrast_std}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_config_multi__d66c76f0f1e5e42f832534b4f9451e5c_qw35sft2_c82d0c91(env, config: dict):
    """Get multiple GIMP configuration settings from gimprc as a dict."""
    config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))")['output'].strip()
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'gimprc not found'}
    keys_to_find = config.get('keys', [])
    settings = {}
    for line in content.decode('utf-8', errors='ignore').splitlines():
        if line.startswith('#') or not line.strip():
            continue
        items = line.strip().lstrip('(').rstrip(')').split()
        if not items:
            continue
        if not keys_to_find or items[0] in keys_to_find:
            settings[items[0]] = ' '.join(items[1:])
    return settings

def get_gimp_dims_brightness__0ffd5d6c48e00cc11530a6062d573f43_qw35sft2_90d64ff9(env, config: dict):
    """Read woman_sitting_by_the_tree.png from the VM and return its dimensions and mean brightness."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/woman_sitting_by_the_tree.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None, 'mean_brightness': None}
    suffix = '.png' if file_path.endswith('.png') else '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        rgb = img.convert('RGB')
        arr = np.array(rgb, dtype=float)
        mean_brightness = float(arr.mean())
        return {'width': img.width, 'height': img.height, 'mean_brightness': mean_brightness}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None, 'mean_brightness': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_mode__a22a6a74169a7c2b7bd4f15c461d97a5_qw35sft2_adb429cd(env, config: dict):
    """Read gate.jpeg from the VM and return its PIL mode and dimensions."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/gate.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None, 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'mode': None, 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_thumbnail__4ddc2b8f762091345213742b2e539174_qw35sft2_33ee6ccc(env, config: dict):
    """Download heron.jpeg from VM and return a small thumbnail pixel array for flip verification."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/heron.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        thumb_w, thumb_h = (64, 43)
        thumb = img.resize((thumb_w, thumb_h), Image.LANCZOS)
        pixels = [list(p) for p in thumb.getdata()]
        return {'width': img.width, 'height': img.height, 'thumbnail_pixels': pixels, 'thumb_size': [thumb_w, thumb_h]}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_session_single_window__940ba8eb8583aa4e66192a1528358acd_qw35sft2_9b24c412(env, config: dict):
    """Read GIMP sessionrc and return single-window-mode state."""
    path_result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/sessionrc'))")
    config_path = path_result.get('output', '').strip() if isinstance(path_result, dict) else ''
    if not config_path:
        return {'single_window_mode': 'unknown'}
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'single_window_mode': 'unknown'}
    content = file_bytes.decode('utf-8', errors='replace')
    for line in content.splitlines():
        stripped = line.strip().lstrip('(').rstrip(')\n')
        parts = stripped.split()
        if parts and parts[0] == 'single-window-mode':
            return {'single_window_mode': parts[-1] if len(parts) > 1 else 'yes'}
    return {'single_window_mode': 'unknown'}

def get_gimp_image_mode__b63a2518d07e9d19e5e5bc7d70fabd81_qw35sft2_019efa2b(env, config: dict):
    """Read the exported PNG from the VM and return its PIL mode."""
    from PIL import Image
    file_path = config.get('path', '/home/user/logo_gray.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None, 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'mode': None, 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_saturation_stats__b674a0dc75a50659a17535e66cb88a78_qw35sft2_a56e2f2c(env, config: dict):
    """Read berries_tv4.png from VM and return avg stddev (contrast) and saturation score."""
    from PIL import Image, ImageStat
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/berries_tv4.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'avg_std': 0.0, 'sat_score': 0.0}
    suffix = '.png'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        stat = ImageStat.Stat(img)
        avg_std = sum(stat.stddev) / 3.0
        arr = np.array(img, dtype=np.float32)
        maxv = arr.max(axis=2)
        minv = arr.min(axis=2)
        sat_score = float(np.where(maxv > 0, (maxv - minv) / maxv, 0.0).mean())
        return {'avg_std': avg_std, 'sat_score': sat_score}
    except Exception as e:
        return {'error': str(e), 'avg_std': 0.0, 'sat_score': 0.0}
    finally:
        os.unlink(tmp_path)

def get_gimp_yellow_centroid__466f916fdda48299f33abc0cff1e195a_qw35sft2_60e5fb31(env, config: dict):
    """Read Triangle_On_The_Side.png from the VM and return the yellow triangle centroid."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/Triangle_On_The_Side.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'centroid_x': None, 'centroid_y': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGBA')
        arr = np.array(img)
        yellow_mask = (arr[:, :, 0] > 150) & (arr[:, :, 1] > 100) & (arr[:, :, 2] < 100) & (arr[:, :, 3] > 50)
        ys, xs = np.where(yellow_mask)
        if len(xs) == 0:
            return {'error': 'No yellow pixels found', 'centroid_x': None, 'centroid_y': None}
        centroid_x = float(xs.mean())
        centroid_y = float(ys.mean())
        return {'centroid_x': centroid_x, 'centroid_y': centroid_y, 'pixel_count': int(len(xs)), 'width': int(arr.shape[1]), 'height': int(arr.shape[0])}
    except Exception as e:
        return {'error': str(e), 'centroid_x': None, 'centroid_y': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_dims__ffabe808aeb093a1184f15812c26092c_qw35sft2_e25cffe9(env, config: dict):
    """Download image file from VM and return its pixel dimensions."""
    import tempfile, os
    path = config.get('path', '/home/user/Desktop/dog_resized.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_and_layer__5e787196d91537456fde84b3a66eeecb_qw35sft2_6a04fddf(env, config: dict):
    """
    Download exported PNG from VM and return:
      - file_w, file_h: full image canvas dimensions
      - bbox_w, bbox_h: non-transparent content bounding box dimensions
    """
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/dog_half.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        file_w, file_h = img.size
        if img.mode == 'RGBA':
            alpha_channel = img.split()[-1]
            bbox = alpha_channel.getbbox()
            if bbox:
                bbox_w = bbox[2] - bbox[0]
                bbox_h = bbox[3] - bbox[1]
            else:
                bbox_w, bbox_h = (0, 0)
        else:
            bbox_w, bbox_h = (file_w, file_h)
        return {'file_w': file_w, 'file_h': file_h, 'bbox_w': bbox_w, 'bbox_h': bbox_h}
    except Exception as e:
        return {'error': str(e), 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_bg_pixel__4d0454ba1a7fe5978c1747f21d797308_qw35sft2_074a772f(env, config: dict):
    """Download exported PNG from VM and return background pixel color at top-left corner."""
    import tempfile
    import os
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL not available', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    path = config.get('path', '/home/user/Desktop/flat_green.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        width, height = img.size
        r, g, b = img.getpixel((10, 10))
        return {'r': int(r), 'g': int(g), 'b': int(b), 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    finally:
        os.unlink(tmp_path)

def get_gimp_hd_png__6d65b36cd0a438b39e8f21308e345a4a_qw35sft2_848c1efa(env, config: dict):
    """Download a PNG from the VM and return its dimensions and file existence."""
    import io
    path = config.get('path', '/home/user/Desktop/hd_image.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(file_bytes))
        return {'width': img.width, 'height': img.height, 'format': img.format}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_gimp_image_size__08c1ae085b68de96a31c7d77e1455141_qw35sft2_87a5475c(env, config: dict):
    """Read export.jpg from the VM and return its pixel dimensions."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/export.jpg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    suffix = os.path.splitext(file_path)[-1] or '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_xcf_layers__ad9b8a9ed8fc4caff86597f320aba5f2_qw35sft2_a7561ace(env, config: dict):
    """Parse XCF file from VM and return layer names and count."""
    import struct
    xcf_path = config.get('path', '/home/user/Desktop/white_background.xcf')
    try:
        file_bytes = env.controller.get_file(xcf_path)
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}
    if not file_bytes:
        return {'error': 'File not found or empty', 'layer_names': [], 'layer_count': 0}
    data = bytes(file_bytes)
    if not data.startswith(b'gimp xcf '):
        return {'error': 'Not a valid XCF file', 'layer_names': [], 'layer_count': 0}
    try:
        null_pos = data.index(b'\x00', 9)
        version_str = data[9:null_pos].decode('ascii', errors='replace')
        pos = null_pos + 1
        if pos + 12 > len(data):
            return {'error': 'File too short', 'layer_names': [], 'layer_count': 0}
        pos += 12
        while pos + 8 <= len(data):
            prop_type = struct.unpack('>I', data[pos:pos + 4])[0]
            prop_size = struct.unpack('>I', data[pos + 4:pos + 8])[0]
            pos += 8 + prop_size
            if prop_type == 0:
                break
        try:
            ver_num = int(version_str.lstrip('v') or 0)
        except Exception:
            ver_num = 0
        offset_size = 8 if ver_num >= 11 else 4
        offset_fmt = '>Q' if ver_num >= 11 else '>I'
        layer_offsets = []
        while pos + offset_size <= len(data):
            offset = struct.unpack(offset_fmt, data[pos:pos + offset_size])[0]
            pos += offset_size
            if offset == 0:
                break
            layer_offsets.append(offset)
        layer_names = []
        for off in layer_offsets:
            if off + 16 > len(data):
                continue
            name_off = off + 12
            name_len = struct.unpack('>I', data[name_off:name_off + 4])[0]
            name_off += 4
            if name_len == 0 or name_off + name_len > len(data):
                continue
            name_bytes = data[name_off:name_off + name_len]
            name = name_bytes.rstrip(b'\x00').decode('utf-8', errors='replace')
            layer_names.append(name)
        return {'layer_names': layer_names, 'layer_count': len(layer_names)}
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}

def get_gimp_xcf_layers__f3ec988a99b0f642c7ac2a56df923954_qw35sft2_59992397(env, config: dict):
    """Stub: superseded by get_gimp_xcf_layers__ad9b8a9ed8fc4caff86597f320aba5f2."""
    return {'layer_names': [], 'layer_count': 0}

def get_gimp_rot180_check__359fadf8cec20093badb6005d3119d0c_qw35sft2_b61f5722(env, config: dict):
    """
    Download original berry.png (still intact on VM) and the exported
    berry_rot180.png. Return pixel-level similarity between the exported
    file and the expected 180-degree rotation of the original.
    """
    import io
    import numpy as np
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL/Pillow not available'}
    orig_bytes = env.controller.get_file('/home/user/Desktop/berry.png')
    if not orig_bytes:
        return {'error': 'Original berry.png not found on VM'}
    result_path = config.get('path', '/home/user/Desktop/berry_rot180.png')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        return {'error': f'Result file not found: {result_path}'}
    try:
        orig_img = Image.open(io.BytesIO(orig_bytes)).convert('RGB')
        result_img = Image.open(io.BytesIO(result_bytes)).convert('RGB')
        expected_img = orig_img.transpose(Image.ROTATE_180)
        if result_img.size != expected_img.size:
            result_img = result_img.resize(expected_img.size, Image.LANCZOS)
        orig_arr = np.array(expected_img).astype(float)
        result_arr = np.array(result_img).astype(float)
        mse = float(np.mean((orig_arr - result_arr) ** 2))
        max_mse = 255.0 ** 2
        similarity = 1.0 - mse / max_mse
        return {'similarity': similarity, 'mse': mse, 'result_size': list(result_img.size), 'expected_size': list(expected_img.size)}
    except Exception as e:
        return {'error': str(e)}

def get_gimp_process__be67663f10a797bb181cd6fa722efa34_qw35sft2_7486e82e(env, config: dict):
    """Check if the GIMP process is currently running on the system."""
    result = env.controller.run_bash_script("pgrep -l gimp 2>/dev/null || echo ''", timeout=10)
    if isinstance(result, dict):
        stdout = result.get('output', result.get('stdout', '')).strip()
    else:
        stdout = str(result).strip()
    running = bool(stdout and 'gimp' in stdout.lower())
    return {'running': running, 'process_info': stdout}

def get_gimp_saturation_size__e189188d82a47045e4921db4c0cbb19b_qw35sft2_1d6c1177(env, config: dict):
    """Download edited image from VM and return mean HSV saturation, width, and height."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    file_bytes = env.controller.get_file('/home/user/Desktop/woman_sitting_by_the_tree2.png')
    if not file_bytes:
        return {'error': 'File not found on VM'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        width, height = img.size
        img_rgb = img.convert('RGB')
        arr = np.array(img_rgb).astype(float) / 255.0
        r, g, b = (arr[:, :, 0], arr[:, :, 1], arr[:, :, 2])
        maxc = np.maximum(np.maximum(r, g), b)
        minc = np.minimum(np.minimum(r, g), b)
        sat = np.where(maxc > 0, (maxc - minc) / maxc, 0.0)
        return {'saturation_mean': float(sat.mean()), 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_grayscale_brightness__72b5517a5a800033bb98cd0fbc3da308_qw35sft2_b713d12c(env, config: dict):
    """Read woman_sitting_by_the_tree.png from the VM and return its PIL mode and mean brightness."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/woman_sitting_by_the_tree.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None, 'mean_brightness': None}
    suffix = '.png' if file_path.endswith('.png') else '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        mode = img.mode
        gray = img.convert('L')
        arr = np.array(gray, dtype=float)
        mean_brightness = float(arr.mean())
        return {'mode': mode, 'mean_brightness': mean_brightness}
    except Exception as e:
        return {'error': str(e), 'mode': None, 'mean_brightness': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_props__e5dd5bd57fb559e638682724445a4fb8_qw35sft2_de8a7ed3(env, config: dict):
    """Get image mode and dimensions from exported PNG file on the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/square_palette_computer.png')
    try:
        from PIL import Image
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'file not found or empty'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            img = Image.open(tmp_path)
            result = {'mode': img.mode, 'width': img.size[0], 'height': img.size[1]}
            img.close()
        finally:
            os.unlink(tmp_path)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_gimp_file_exists__6e337deb47fb4ded97926f2aaf4149ad_qw35sft2_4ca91268(env, config: dict):
    """Retrieve the exported BMP file and return header bytes for format validation."""
    file_path = config.get('path', '/home/user/logo.bmp')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is not None and len(file_bytes) > 0:
            return {'exists': True, 'size': len(file_bytes), 'header_bytes': file_bytes[:2]}
        return {'exists': False, 'size': 0, 'header_bytes': b''}
    except Exception as e:
        return {'exists': False, 'error': str(e), 'header_bytes': b''}

def get_image_brightness_stat__1dd8caa614c12c2ce7c1b58be9a02c11_qw35sft2_dce61cdc(env, config: dict):
    """Get average RGB brightness of an image file from the VM."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    vm_path = config.get('path', '')
    if not vm_path:
        return {'error': 'No path specified in config'}
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': f'File not found or empty: {vm_path}', 'file_found': False}
    suffix = os.path.splitext(vm_path)[1] or '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        arr = np.array(img, dtype=float)
        avg_brightness = float(arr.mean())
        return {'avg_brightness': avg_brightness, 'file_found': True}
    except Exception as e:
        return {'error': str(e), 'file_found': False}
    finally:
        os.unlink(tmp_path)

def get_gimp_jpeg_stats__99df0d1fa420bfe8f4d58bfec01f5388_qw35sft2_9f161e6b(env, config: dict):
    """Read berries_tv2.jpg from VM and return its PIL format and avg stddev (contrast proxy)."""
    from PIL import Image, ImageStat
    file_path = config.get('path', '/home/user/Desktop/berries_tv2.jpg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'format': None, 'avg_std': 0.0}
    suffix = '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_format = img.format
        img_rgb = img.convert('RGB')
        stat = ImageStat.Stat(img_rgb)
        avg_std = sum(stat.stddev) / 3.0
        return {'format': img_format, 'avg_std': avg_std}
    except Exception as e:
        return {'error': str(e), 'format': None, 'avg_std': 0.0}
    finally:
        os.unlink(tmp_path)

def get_image_dimensions__1649a3c87969acc3b9b0ad261438802f_qw35sft2_8ed03cf1(env, config: dict):
    """Get width and height of the PNG image on the Desktop using PIL."""
    import tempfile, os
    path = config.get('path', '/home/user/Desktop/dog_with_background.png')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'File not found', 'width': None, 'height': None}
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(file_bytes))
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_gimp_image_dims__2254c4ac73b2e0e2d0bc95e32c93006c_qw35sft2_ae629eeb(env, config: dict):
    """Read gate.jpeg from the VM and return its dimensions."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/gate.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_grayscale_check__c8310c6f68cacbb7f806947f886f9156_qw35sft2_75ba4ee6(env, config: dict):
    """Download exported grayscale image and return its color mode and dimensions."""
    import tempfile, os
    path = config.get('path', '/home/user/Desktop/dog_grayscale.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'mode': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_mode__5cf8f74046cc654cd234f12e0bd56fa2_qw35sft2_e1f7e6c9(env, config: dict):
    """Download heron.jpeg from VM and return its PIL mode and dimensions."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/heron.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None, 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'mode': None, 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_icon_theme__8b7e0a81347ead101665acba492eaaa3_qw35sft2_75b7fca1(env, config: dict):
    """Read GIMP gimprc and return the current icon-theme setting."""
    path_result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))")
    config_path = path_result.get('output', '').strip() if isinstance(path_result, dict) else ''
    if not config_path:
        return {'icon_theme': 'unknown'}
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'icon_theme': 'unknown'}
    content = file_bytes.decode('utf-8', errors='replace')
    for line in content.splitlines():
        stripped = line.strip().lstrip('(').rstrip(')\n')
        parts = stripped.split()
        if parts and parts[0] == 'icon-theme':
            icon_val = ' '.join(parts[1:]).strip('"\'') if len(parts) > 1 else 'unknown'
            return {'icon_theme': icon_val}
    return {'icon_theme': 'unknown'}

def get_gimp_png_orientation__65c7fe1c7d5d90a17513cb0cbcf64984_qw35sft2_38622937(env, config: dict):
    """Download the exported PNG and return its dimensions and orientation flag."""
    import tempfile, os
    vm_path = config.get('path', '/home/user/Desktop/orange_background.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'File not found: ' + vm_path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        width, height = img.size
        return {'width': width, 'height': height, 'is_portrait': height > width}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_layer_bbox__6fac91d7fdc9ac567c2522f5862ea0c3_qw35sft2_ce514f78(env, config: dict):
    """Download exported PNG from VM and return file dims + non-transparent bounding box dims."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/dog_256.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        file_w, file_h = img.size
        if img.mode == 'RGBA':
            alpha_channel = img.split()[-1]
            bbox = alpha_channel.getbbox()
            if bbox:
                bbox_w = bbox[2] - bbox[0]
                bbox_h = bbox[3] - bbox[1]
            else:
                bbox_w, bbox_h = (0, 0)
        else:
            bbox_w, bbox_h = (file_w, file_h)
        return {'file_w': file_w, 'file_h': file_h, 'bbox_w': bbox_w, 'bbox_h': bbox_h}
    except Exception as e:
        return {'error': str(e), 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_exported_png__2034cd2081886ce4d15c41eaddab9687_qw35sft2_569c72d7(env, config: dict):
    """Read the exported PNG from the VM and return its dimensions and mode."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/output.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None, 'mode': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'mode': img.mode}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None, 'mode': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_bg_pixel__c69df4a705c1504ceebda9205ba51ae8_qw35sft2_123f6c98(env, config: dict):
    """Download exported PNG from VM and return background pixel color at top-left corner."""
    import tempfile
    import os
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL not available', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    path = config.get('path', '/home/user/Desktop/verify_green_v0.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        width, height = img.size
        r, g, b = img.getpixel((10, 10))
        return {'r': int(r), 'g': int(g), 'b': int(b), 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    finally:
        os.unlink(tmp_path)

def get_gimp_hflip_check__e47abffb5b090802d85e5308a4df8da8_qw35sft2_689caf68(env, config: dict):
    """
    Download original berry.png (still intact on VM) and the exported
    berry_flipped.png. Return pixel-level similarity between the exported
    file and the expected horizontal-flip of the original.
    """
    import io
    import numpy as np
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL/Pillow not available'}
    orig_bytes = env.controller.get_file('/home/user/Desktop/berry.png')
    if not orig_bytes:
        return {'error': 'Original berry.png not found on VM'}
    result_path = config.get('path', '/home/user/Desktop/berry_flipped.png')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        return {'error': f'Result file not found: {result_path}'}
    try:
        orig_img = Image.open(io.BytesIO(orig_bytes)).convert('RGB')
        result_img = Image.open(io.BytesIO(result_bytes)).convert('RGB')
        expected_img = orig_img.transpose(Image.FLIP_LEFT_RIGHT)
        if result_img.size != expected_img.size:
            result_img = result_img.resize(expected_img.size, Image.LANCZOS)
        orig_arr = np.array(expected_img).astype(float)
        result_arr = np.array(result_img).astype(float)
        mse = float(np.mean((orig_arr - result_arr) ** 2))
        max_mse = 255.0 ** 2
        similarity = 1.0 - mse / max_mse
        return {'similarity': similarity, 'mse': mse, 'result_size': list(result_img.size), 'expected_size': list(expected_img.size)}
    except Exception as e:
        return {'error': str(e)}

def get_gimp_theme_setting__a2417c567c82dd8e5e291c848fcf86af_qw35sft2_31a64438(env, config: dict):
    """Read gimprc and return the current theme setting."""
    path = config.get('path', '/home/user/.config/GIMP/2.10/gimprc')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'gimprc not found', 'theme': None}
    content = file_bytes.decode('utf-8', errors='replace')
    theme = None
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith('(theme '):
            val = stripped[7:].rstrip(')')
            theme = val.strip().strip('"').lower()
    return {'theme': theme}

def get_gimp_saturation_stats__6b90b70b3d1b161e089c6f4f8582392c_qw35sft2_ad606690(env, config: dict):
    """Download the edited image from VM and compute mean HSV saturation."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    file_bytes = env.controller.get_file('/home/user/Desktop/woman_sitting_by_the_tree2.png')
    if not file_bytes:
        return {'error': 'File not found on VM'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        arr = np.array(img).astype(float) / 255.0
        r, g, b = (arr[:, :, 0], arr[:, :, 1], arr[:, :, 2])
        maxc = np.maximum(np.maximum(r, g), b)
        minc = np.minimum(np.minimum(r, g), b)
        sat = np.where(maxc > 0, (maxc - minc) / maxc, 0.0)
        return {'saturation_mean': float(sat.mean())}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_canvas_and_layer__0d4c5e593127b30f4279483bc536e453_qw35sft2_6c86c613(env, config: dict):
    """Check for canvas.png on Desktop (exact match) and 'Square' layer in saved XCF."""
    result_ls = env.controller.run_bash_script('ls /home/user/Desktop/', timeout=10)
    listing = result_ls.get('output', '') if isinstance(result_ls, dict) else str(result_ls)
    files = [f.strip() for f in listing.strip().split('\n') if f.strip()]
    png_exists = 'canvas.png' in files
    xcf_bytes = env.controller.get_file('/home/user/Desktop/white_background.xcf')
    square_layer_found = False
    if xcf_bytes:
        square_layer_found = b'Square\x00' in xcf_bytes
    return {'png_exists': png_exists, 'square_layer_found': square_layer_found}

def get_desktop_jpeg_exists__e78c01185325a38b78f35c525682fa2c_qw35sft2_33aee3e6(env, config: dict):
    """Check if yicun.jpg exists on the Desktop and is a valid JPEG file."""
    jpeg_path = '/home/user/Desktop/yicun.jpg'
    result = env.controller.run_bash_script(f"test -f '{jpeg_path}' && echo 'exists' || echo 'missing'", timeout=10)
    if isinstance(result, dict):
        stdout = result.get('output', result.get('stdout', '')).strip()
    else:
        stdout = str(result).strip()
    if stdout != 'exists':
        return {'exists': False, 'valid_jpeg': False, 'path': jpeg_path}
    try:
        file_bytes = env.controller.get_file(jpeg_path)
        if file_bytes and len(file_bytes) > 2 and (file_bytes[:2] == b'\xff\xd8'):
            return {'exists': True, 'valid_jpeg': True, 'path': jpeg_path}
        else:
            return {'exists': True, 'valid_jpeg': False, 'path': jpeg_path}
    except Exception:
        return {'exists': True, 'valid_jpeg': False, 'path': jpeg_path}

def get_gimp_image_color_info__4728404e320da0722e3f0bbf27779384_qw35sft2_7a1c2b28(env, config: dict):
    """Get image mode and unique color count from an indexed PNG on the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/palette128_computer.png')
    try:
        from PIL import Image
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'file not found or empty'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            img = Image.open(tmp_path)
            mode = img.mode
            if mode == 'P':
                color_count = len(set(img.getdata()))
            else:
                color_count = -1
            result = {'mode': mode, 'color_count': color_count}
            img.close()
        finally:
            os.unlink(tmp_path)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_gimp_jpeg_export__3494183243a5e3cf64667682e9d5cf69_qw35sft2_b423750f(env, config: dict):
    """
    Check if dark_photo.jpg exists on the Desktop and return its mean RGB brightness.
    """
    from PIL import Image
    import numpy as np
    jpeg_path = config.get('jpeg_path', '/home/user/Desktop/dark_photo.jpg')
    result = {'jpeg_exists': False, 'jpeg_mean_brightness': None}
    jpeg_bytes = env.controller.get_file(jpeg_path)
    if jpeg_bytes:
        result['jpeg_exists'] = True
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(jpeg_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path).convert('RGB')
            arr = np.array(img, dtype=float)
            result['jpeg_mean_brightness'] = float(arr.mean())
        except Exception as e:
            result['jpeg_error'] = str(e)
        finally:
            os.unlink(tmp_path)
    return result

def get_gimp_gray_stats__a80023f11b098859420cd69b0067e271_qw35sft2_e03c66fa(env, config: dict):
    """Read berries_tv3.png from VM and return its PIL mode and grayscale stddev."""
    from PIL import Image, ImageStat
    file_path = config.get('path', '/home/user/Desktop/berries_tv3.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'mode': None, 'gray_std': 0.0}
    suffix = '.png'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        mode = img.mode
        img_gray = img.convert('L')
        stat = ImageStat.Stat(img_gray)
        gray_std = stat.stddev[0]
        return {'mode': mode, 'gray_std': gray_std}
    except Exception as e:
        return {'error': str(e), 'mode': None, 'gray_std': 0.0}
    finally:
        os.unlink(tmp_path)

def get_desktop_xcf__30619c16638f3c14c2b868710d3b511a_qw35sft2_3341f65e(env, config: dict):
    """List XCF files on Desktop and read canvas dimensions from the target XCF header."""
    ls_result = env.controller.run_bash_script('ls ~/Desktop/ 2>/dev/null', timeout=10)
    if isinstance(ls_result, dict):
        stdout = ls_result.get('output', ls_result.get('stdout', '')).strip()
    else:
        stdout = str(ls_result).strip()
    files = [f.strip() for f in stdout.split('\n') if f.strip()] if stdout else []
    filename = config.get('filename', 'new_image.xcf')
    width = None
    height = None
    if filename in files:
        dim_script = "python3 << 'PYEOF'\nimport struct, os\ntry:\n    path = os.path.expanduser('~/Desktop/" + filename + "')\n    with open(path, 'rb') as f:\n        data = f.read(30)\n    null_pos = data.index(b'\\x00', 9)\n    w, h = struct.unpack('>II', data[null_pos+1:null_pos+9])\n    print(w, h)\nexcept Exception:\n    print('error')\nPYEOF\n"
        dim_result = env.controller.run_bash_script(dim_script, timeout=10)
        if isinstance(dim_result, dict):
            dim_out = dim_result.get('output', dim_result.get('stdout', '')).strip()
        else:
            dim_out = str(dim_result).strip()
        parts = dim_out.split()
        if len(parts) == 2 and parts[0] != 'error':
            try:
                width = int(parts[0])
                height = int(parts[1])
            except ValueError:
                pass
    return {'files': files, 'width': width, 'height': height}

def get_gimp_config_multi__59d11816792f7732bb8631499d5a3dca_qw35sft2_434bfb14(env, config: dict):
    """Get multiple GIMP configuration settings from gimprc as a dict."""
    config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))")['output'].strip()
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'gimprc not found'}
    keys_to_find = config.get('keys', [])
    settings = {}
    for line in content.decode('utf-8', errors='ignore').splitlines():
        if line.startswith('#') or not line.strip():
            continue
        items = line.strip().lstrip('(').rstrip(')').split()
        if not items:
            continue
        if not keys_to_find or items[0] in keys_to_find:
            settings[items[0]] = ' '.join(items[1:])
    return settings

def get_gimp_flipped_file__f5b72d9abaa33261d0da153e643d90eb_qw35sft2_02e67ab3(env, config: dict):
    """Download flipped and original images from VM, compare pixel arrays to verify horizontal flip."""
    import tempfile, os
    import numpy as np
    from PIL import Image, ImageOps
    flipped_path = config.get('path', '/home/user/Desktop/dog_flipped.png')
    original_path = config.get('original_path', '/home/user/Desktop/dog_with_background.png')
    flipped_bytes = env.controller.get_file(flipped_path)
    if not flipped_bytes:
        return {'error': 'Exported file not found', 'exists': False}
    original_bytes = env.controller.get_file(original_path)
    if not original_bytes:
        return {'error': 'Original file not found', 'exists': False}
    tmp_flipped = None
    tmp_original = None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(flipped_bytes)
            tmp_flipped = f.name
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(original_bytes)
            tmp_original = f.name
        flipped_img = Image.open(tmp_flipped).convert('RGB')
        original_img = Image.open(tmp_original).convert('RGB')
        expected_flipped = ImageOps.mirror(original_img)
        flipped_arr = np.array(flipped_img)
        expected_arr = np.array(expected_flipped)
        flip_match = bool(np.array_equal(flipped_arr, expected_arr))
        return {'exists': True, 'width': flipped_img.width, 'height': flipped_img.height, 'flip_match': flip_match}
    except Exception as e:
        return {'error': str(e), 'exists': False}
    finally:
        if tmp_flipped and os.path.exists(tmp_flipped):
            os.unlink(tmp_flipped)
        if tmp_original and os.path.exists(tmp_original):
            os.unlink(tmp_original)

def get_gimp_image_flip__917488a366fb0a5d7f8700214812c742_qw35sft2_ea5a8b0d(env, config: dict):
    """Read gate.jpeg from the VM and compute left/right half brightness averages for flip detection.

    Original image (2850x5070) statistics:
      left_half_avg  ≈ 109.72
      right_half_avg ≈  77.94

    After a horizontal flip these values swap:
      left_half_avg  ≈  77.94
      right_half_avg ≈ 109.72
    """
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/gate.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'left_half_avg': None, 'right_half_avg': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        arr = np.array(img, dtype=float)
        mid = arr.shape[1] // 2
        left_avg = float(arr[:, :mid, :].mean())
        right_avg = float(arr[:, mid:, :].mean())
        return {'left_half_avg': left_avg, 'right_half_avg': right_avg, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'left_half_avg': None, 'right_half_avg': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_rotated_dims__480c196a41fcbd8c06328eee76532eca_qw35sft2_a1188a1a(env, config: dict):
    """Download heron.jpeg from VM and return its dimensions after rotation."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/heron.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'mode': img.mode}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_file_exists__9dede0efb5d54dc2e4df137a06f2c4ba_qw35sft2_e1eea9d6(env, config: dict):
    """Check if the specified file exists on the VM by attempting to retrieve it."""
    file_path = config.get('path', '/home/user/Documents/export.jpg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'size': 0}
    return {'exists': True, 'size': len(file_bytes)}

def get_gimp_dark_theme__48287b5138535fb1248d29dc8ceabce5_qw35sft2_c404d8a1(env, config: dict):
    """Read GIMP gimprc and return the current theme setting."""
    path_result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))")
    config_path = path_result.get('output', '').strip() if isinstance(path_result, dict) else ''
    if not config_path:
        return {'theme': 'unknown'}
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'theme': 'unknown'}
    content = file_bytes.decode('utf-8', errors='replace')
    for line in content.splitlines():
        stripped = line.strip().lstrip('(').rstrip(')\n')
        parts = stripped.split()
        if parts and parts[0] == 'theme':
            theme_val = ' '.join(parts[1:]).strip('"\'') if len(parts) > 1 else 'unknown'
            return {'theme': theme_val}
    return {'theme': 'unknown'}

def get_image_dimensions__9d59aa3517eac66c5b7e40985af8d11e_qw35sft2_82cd2c39(env, config: dict):
    """Get width and height of the PNG image on the Desktop using PIL."""
    path = config.get('path', '/home/user/Desktop/dog_with_background.png')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'File not found', 'width': None, 'height': None}
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(file_bytes))
        return {'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}

def get_gimp_png_text_pos__f282f693e16d842a8fd79d730244ff72_qw35sft2_f5a11855(env, config: dict):
    """Download the exported PNG and compute the horizontal center of dark (text) pixels."""
    import tempfile, os
    vm_path = config.get('path', '/home/user/Desktop/orange_background.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'File not found: ' + vm_path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        import numpy as np
        img = Image.open(tmp_path).convert('RGB')
        img_arr = np.array(img)
        width = img_arr.shape[1]
        r = img_arr[:, :, 0].astype(int)
        g = img_arr[:, :, 1].astype(int)
        b = img_arr[:, :, 2].astype(int)
        is_dark = (r < 80) & (g < 80) & (b < 80)
        dark_count = int(is_dark.sum())
        if dark_count == 0:
            return {'error': 'No dark/text pixels found', 'width': width}
        dark_x = np.where(is_dark)[1]
        text_center_x = int(dark_x.mean())
        return {'width': width, 'text_center_x': text_center_x, 'is_on_left_half': text_center_x < width // 2}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_image_brightness_stat__53b38a939b32dab621a5954331033a3c_qw35sft2_3c0a82cc(env, config: dict):
    """Get average RGB brightness of an image file from the VM."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    vm_path = config.get('path', '')
    if not vm_path:
        return {'error': 'No path specified in config'}
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': f'File not found or empty: {vm_path}', 'file_found': False}
    suffix = os.path.splitext(vm_path)[1] or '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        arr = np.array(img, dtype=float)
        avg_brightness = float(arr.mean())
        return {'avg_brightness': avg_brightness, 'file_found': True}
    except Exception as e:
        return {'error': str(e), 'file_found': False}
    finally:
        os.unlink(tmp_path)

def get_gimp_yellow_centroid__569dedcfdc3e7d291a20f0e6aca641c0_qw35sft2_8d4fd89e(env, config: dict):
    """Read Triangle_On_The_Side.png from the VM and return the yellow triangle centroid."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/Triangle_On_The_Side.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'centroid_x': None, 'centroid_y': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGBA')
        arr = np.array(img)
        yellow_mask = (arr[:, :, 0] > 150) & (arr[:, :, 1] > 100) & (arr[:, :, 2] < 100) & (arr[:, :, 3] > 50)
        ys, xs = np.where(yellow_mask)
        if len(xs) == 0:
            return {'error': 'No yellow pixels found', 'centroid_x': None, 'centroid_y': None}
        centroid_x = float(xs.mean())
        centroid_y = float(ys.mean())
        return {'centroid_x': centroid_x, 'centroid_y': centroid_y, 'pixel_count': int(len(xs)), 'width': int(arr.shape[1]), 'height': int(arr.shape[0])}
    except Exception as e:
        return {'error': str(e), 'centroid_x': None, 'centroid_y': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_dims__2bd82e314cabe351963c3af4824b9b87_qw35sft2_c504650a(env, config: dict):
    """Get image dimensions from the scaled PNG file on the VM."""
    from PIL import Image
    file_path = config.get('path', '/home/user/logo_small.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'mode': img.mode}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_layer_bbox_opacity__c6fafc0c96511b4f7b94b1c390c2f8bd_qw35sft2_a1553309(env, config: dict):
    """Download exported PNG from VM and return file dims, non-transparent bbox dims, and max alpha value."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/dog_opacity.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None, 'max_alpha': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        file_w, file_h = img.size
        if img.mode == 'RGBA':
            alpha_channel = img.split()[-1]
            bbox = alpha_channel.getbbox()
            if bbox:
                bbox_w = bbox[2] - bbox[0]
                bbox_h = bbox[3] - bbox[1]
            else:
                bbox_w, bbox_h = (0, 0)
            alpha_min, alpha_max = alpha_channel.getextrema()
            max_alpha = alpha_max
        else:
            bbox_w, bbox_h = (file_w, file_h)
            max_alpha = 255
        return {'file_w': file_w, 'file_h': file_h, 'bbox_w': bbox_w, 'bbox_h': bbox_h, 'max_alpha': max_alpha}
    except Exception as e:
        return {'error': str(e), 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None, 'max_alpha': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_bg_pixel__3cdf6e4d2fe41e5dc87ba95dd55c83c0_qw35sft2_b80b5fe7(env, config: dict):
    """Download exported PNG from VM and return background pixel color at top-left corner."""
    import tempfile
    import os
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL not available', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    path = config.get('path', '/home/user/Desktop/verify_red_v1.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        width, height = img.size
        r, g, b = img.getpixel((10, 10))
        return {'r': int(r), 'g': int(g), 'b': int(b), 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    finally:
        os.unlink(tmp_path)

def get_gimp_vflip_check__e3f424dc6a1401f7936bf6fab235a0b4_qw35sft2_f99490de(env, config: dict):
    """
    Download original berry.png (still intact on VM) and the exported
    berry_vflip.png. Return pixel-level similarity between the exported
    file and the expected vertical-flip of the original.
    """
    import io
    import numpy as np
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL/Pillow not available'}
    orig_bytes = env.controller.get_file('/home/user/Desktop/berry.png')
    if not orig_bytes:
        return {'error': 'Original berry.png not found on VM'}
    result_path = config.get('path', '/home/user/Desktop/berry_vflip.png')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        return {'error': f'Result file not found: {result_path}'}
    try:
        orig_img = Image.open(io.BytesIO(orig_bytes)).convert('RGB')
        result_img = Image.open(io.BytesIO(result_bytes)).convert('RGB')
        expected_img = orig_img.transpose(Image.FLIP_TOP_BOTTOM)
        if result_img.size != expected_img.size:
            result_img = result_img.resize(expected_img.size, Image.LANCZOS)
        orig_arr = np.array(expected_img).astype(float)
        result_arr = np.array(result_img).astype(float)
        mse = float(np.mean((orig_arr - result_arr) ** 2))
        max_mse = 255.0 ** 2
        similarity = 1.0 - mse / max_mse
        return {'similarity': similarity, 'mse': mse, 'result_size': list(result_img.size), 'expected_size': list(expected_img.size)}
    except Exception as e:
        return {'error': str(e)}

def get_gimp_icon_theme__a072fad36e1102073ae062e8b644b246_qw35sft2_6918c39c(env, config: dict):
    """Read gimprc and return the current icon-theme setting."""
    path = config.get('path', '/home/user/.config/GIMP/2.10/gimprc')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'gimprc not found', 'icon_theme': None}
    content = file_bytes.decode('utf-8', errors='replace')
    icon_theme = None
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith('(icon-theme '):
            val = stripped[12:].rstrip(')')
            icon_theme = val.strip().strip('"').lower()
    return {'icon_theme': icon_theme}

def get_gimp_image_props__94a488f0af5ba53f580224747ebe24fe_qw35sft2_3884b707(env, config: dict):
    """Get image mode and dimensions from exported PNG file on the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/grayscale_scaled_computer.png')
    try:
        from PIL import Image
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'file not found or empty'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            img = Image.open(tmp_path)
            result = {'mode': img.mode, 'width': img.size[0], 'height': img.size[1]}
            img.close()
        finally:
            os.unlink(tmp_path)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_gimp_xcf_layers__5ffee09eda4febddd3875d5ca422ec9b_qw35sft2_0ad9a860(env, config: dict):
    """Parse XCF file from VM and return layer names and count."""
    import struct
    xcf_path = config.get('path', '/home/user/Desktop/white_background.xcf')
    try:
        file_bytes = env.controller.get_file(xcf_path)
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}
    if not file_bytes:
        return {'error': 'File not found or empty', 'layer_names': [], 'layer_count': 0}
    data = bytes(file_bytes)
    if not data.startswith(b'gimp xcf '):
        return {'error': 'Not a valid XCF file', 'layer_names': [], 'layer_count': 0}
    try:
        null_pos = data.index(b'\x00', 9)
        version_str = data[9:null_pos].decode('ascii', errors='replace')
        pos = null_pos + 1
        if pos + 12 > len(data):
            return {'error': 'File too short', 'layer_names': [], 'layer_count': 0}
        pos += 12
        while pos + 8 <= len(data):
            prop_type = struct.unpack('>I', data[pos:pos + 4])[0]
            prop_size = struct.unpack('>I', data[pos + 4:pos + 8])[0]
            pos += 8 + prop_size
            if prop_type == 0:
                break
        try:
            ver_num = int(version_str.lstrip('v') or 0)
        except Exception:
            ver_num = 0
        offset_size = 8 if ver_num >= 11 else 4
        offset_fmt = '>Q' if ver_num >= 11 else '>I'
        layer_offsets = []
        while pos + offset_size <= len(data):
            offset = struct.unpack(offset_fmt, data[pos:pos + offset_size])[0]
            pos += offset_size
            if offset == 0:
                break
            layer_offsets.append(offset)
        layer_names = []
        for off in layer_offsets:
            if off + 16 > len(data):
                continue
            name_off = off + 12
            name_len = struct.unpack('>I', data[name_off:name_off + 4])[0]
            name_off += 4
            if name_len == 0 or name_off + name_len > len(data):
                continue
            name_bytes = data[name_off:name_off + name_len]
            name = name_bytes.rstrip(b'\x00').decode('utf-8', errors='replace')
            layer_names.append(name)
        return {'layer_names': layer_names, 'layer_count': len(layer_names)}
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}

def get_gimp_dims_brightness__1abf0259addada4e9320a412de93a720_qw35sft2_f0440aa5(env, config: dict):
    """Read woman_sitting_by_the_tree.png from the VM and return its dimensions and mean brightness."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/woman_sitting_by_the_tree.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None, 'mean_brightness': None}
    suffix = '.png' if file_path.endswith('.png') else '.jpg'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        rgb = img.convert('RGB')
        arr = np.array(rgb, dtype=float)
        mean_brightness = float(arr.mean())
        return {'width': img.width, 'height': img.height, 'mean_brightness': mean_brightness}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None, 'mean_brightness': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_config_multi__d46e54f2b8fa1eafde00e6af910958bf_qw35sft2_3aa5062d(env, config: dict):
    """Get multiple GIMP configuration settings from gimprc as a dict."""
    config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/gimprc'))")['output'].strip()
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'gimprc not found'}
    keys_to_find = config.get('keys', [])
    settings = {}
    for line in content.decode('utf-8', errors='ignore').splitlines():
        if line.startswith('#') or not line.strip():
            continue
        items = line.strip().lstrip('(').rstrip(')').split()
        if not items:
            continue
        if not keys_to_find or items[0] in keys_to_find:
            settings[items[0]] = ' '.join(items[1:])
    return settings

def get_gimp_theme_setting__ff19a7440acd5edb92781b332298a214_qw35sft2_f3f4120e(env, config: dict):
    """Read GIMP gimprc configuration file and return its content for theme checking."""
    gimprc_path = '/home/user/.config/GIMP/2.10/gimprc'
    file_bytes = env.controller.get_file(gimprc_path)
    if not file_bytes:
        return {'error': 'gimprc not found', 'config': ''}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        return {'error': str(e), 'config': ''}
    return {'config': content}

def get_gimp_scale_stats__4ce9d4657814d8ed93df9fa8b61db4cf_qw35sft2_8a4fa70a(env, config: dict):
    """Read berries_tv1.png from VM and return width, height, and avg stddev (contrast proxy)."""
    from PIL import Image, ImageStat
    file_path = config.get('path', '/home/user/Desktop/berries_tv1.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': 0, 'height': 0, 'avg_std': 0.0}
    suffix = '.png'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_rgb = img.convert('RGB')
        stat = ImageStat.Stat(img_rgb)
        avg_std = sum(stat.stddev) / 3.0
        return {'width': img.width, 'height': img.height, 'avg_std': avg_std}
    except Exception as e:
        return {'error': str(e), 'width': 0, 'height': 0, 'avg_std': 0.0}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_export__fcf5b1c8c41f94f2ef23ac5bc8b1ca54_qw35sft2_ebcdfd42(env, config: dict):
    """Check whether a PNG file was exported to the Desktop in GIMP.

    Looks for /home/user/Desktop/gate.png first; if absent, checks for any
    .png file on the Desktop so that alternate filenames are still detected.
    """
    desktop = config.get('desktop_path', '/home/user/Desktop')
    primary_path = os.path.join(desktop, 'gate.png')
    result = env.controller.run_bash_script(f"test -f '{primary_path}' && echo 'found' || echo 'not_found'", timeout=10)
    output = ''
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', ''))
    else:
        output = str(result)
    if 'found' in output and 'not_found' not in output:
        return {'has_png': True, 'path': primary_path}
    result2 = env.controller.run_bash_script(f"ls '{desktop}'/*.png 2>/dev/null | head -5", timeout=10)
    output2 = ''
    if isinstance(result2, dict):
        output2 = result2.get('output', result2.get('stdout', ''))
    else:
        output2 = str(result2)
    png_files = [line.strip() for line in output2.splitlines() if line.strip().endswith('.png')]
    return {'has_png': len(png_files) > 0, 'path': png_files[0] if png_files else None}

def get_gimp_saturation_sharpness__cc63708ed8a0e82770a04258c8e27328_qw35sft2_c9a74c7a(env, config: dict):
    """Download edited image and return mean HSV saturation and Laplacian-variance sharpness."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
        from scipy.ndimage import laplace
    except ImportError:
        return {'error': 'PIL, numpy, or scipy not available'}
    file_bytes = env.controller.get_file('/home/user/Desktop/woman_sitting_by_the_tree2.png')
    if not file_bytes:
        return {'error': 'File not found on VM'}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_rgb = img.convert('RGB')
        arr = np.array(img_rgb).astype(float) / 255.0
        r, g, b = (arr[:, :, 0], arr[:, :, 1], arr[:, :, 2])
        maxc = np.maximum(np.maximum(r, g), b)
        minc = np.minimum(np.minimum(r, g), b)
        sat = np.where(maxc > 0, (maxc - minc) / maxc, 0.0)
        img_gray = img.convert('L')
        gray_arr = np.array(img_gray).astype(float)
        laplacian_var = float(laplace(gray_arr).var())
        return {'saturation_mean': float(sat.mean()), 'laplacian_var': laplacian_var}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_docks_window__e872a1deec6ced899c5abf40a9c68092_qw35sft2_42f2e145(env, config: dict):
    """Read GIMP sessionrc and return hide-docks and single-window-mode states."""
    path_result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/GIMP/2.10/sessionrc'))")
    config_path = path_result.get('output', '').strip() if isinstance(path_result, dict) else ''
    if not config_path:
        return {'hide_docks': 'no', 'single_window_mode': 'unknown'}
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'hide_docks': 'no', 'single_window_mode': 'unknown'}
    content = file_bytes.decode('utf-8', errors='replace')
    hide_docks = 'no'
    single_window_mode = 'unknown'
    for line in content.splitlines():
        stripped = line.strip().lstrip('(').rstrip(')\n')
        parts = stripped.split()
        if not parts:
            continue
        if parts[0] == 'hide-docks':
            hide_docks = parts[-1] if len(parts) > 1 else 'yes'
        elif parts[0] == 'single-window-mode':
            single_window_mode = parts[-1] if len(parts) > 1 else 'yes'
    return {'hide_docks': hide_docks, 'single_window_mode': single_window_mode}

def get_gimp_png_square__ba207a15040c6bc16b5ad1659309ab17_qw35sft2_f3c5241b(env, config: dict):
    """Download the exported PNG and check whether it is square (width == height)."""
    import tempfile, os
    vm_path = config.get('path', '/home/user/Desktop/orange_background.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'File not found: ' + vm_path}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        width, height = img.size
        return {'width': width, 'height': height, 'is_square': width == height}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_yellow_centroid__bc6d20fd4fedacbbf14f4629287f2cc1_qw35sft2_56a18ff2(env, config: dict):
    """Read Triangle_On_The_Side.png from the VM and return the yellow triangle centroid."""
    from PIL import Image
    import numpy as np
    file_path = config.get('path', '/home/user/Desktop/Triangle_On_The_Side.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'centroid_x': None, 'centroid_y': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGBA')
        arr = np.array(img)
        yellow_mask = (arr[:, :, 0] > 150) & (arr[:, :, 1] > 100) & (arr[:, :, 2] < 100) & (arr[:, :, 3] > 50)
        ys, xs = np.where(yellow_mask)
        if len(xs) == 0:
            return {'error': 'No yellow pixels found', 'centroid_x': None, 'centroid_y': None}
        centroid_x = float(xs.mean())
        centroid_y = float(ys.mean())
        return {'centroid_x': centroid_x, 'centroid_y': centroid_y, 'pixel_count': int(len(xs)), 'width': int(arr.shape[1]), 'height': int(arr.shape[0])}
    except Exception as e:
        return {'error': str(e), 'centroid_x': None, 'centroid_y': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_file_size__5c243e9c006808209adf482ad7ac1fe7_qw35sft2_cadec8c0(env, config: dict):
    """Read export.jpg from the VM and return its file size in bytes."""
    file_path = config.get('path', '/home/user/Desktop/export.jpg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'size': None}
    return {'size': len(file_bytes)}

def get_jpeg_file_exists__132eea53d5785c870b331091d6b9fa18_qw35sft2_b2f768ea(env, config: dict):
    """Check whether a JPEG file was exported to the Desktop by the agent."""
    try:
        result = env.controller.run_bash_script('ls /home/user/Desktop/*.jpg /home/user/Desktop/*.jpeg 2>/dev/null | head -5', timeout=15)
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        lines = [ln.strip() for ln in output.splitlines() if ln.strip()]
        return {'jpeg_exists': len(lines) > 0, 'found_files': lines}
    except Exception as e:
        return {'jpeg_exists': False, 'error': str(e)}

def get_gimp_image_dims__392ddc5ce54a4061c2326df12808cf9c_qw35sft2_21889611(env, config: dict):
    """Download heron.jpeg from VM and return its dimensions."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/heron.jpeg')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'width': img.width, 'height': img.height, 'mode': img.mode}
    except Exception as e:
        return {'error': str(e), 'width': None, 'height': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_layer_bbox__4cd9609153e6f966f651b58aa2377fc7_qw35sft2_58886e85(env, config: dict):
    """Download exported PNG from VM and return file dims + non-transparent bounding box dims."""
    from PIL import Image
    file_path = config.get('path', '/home/user/Desktop/dog_800.png')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        file_w, file_h = img.size
        if img.mode == 'RGBA':
            alpha_channel = img.split()[-1]
            bbox = alpha_channel.getbbox()
            if bbox:
                bbox_w = bbox[2] - bbox[0]
                bbox_h = bbox[3] - bbox[1]
            else:
                bbox_w, bbox_h = (0, 0)
        else:
            bbox_w, bbox_h = (file_w, file_h)
        return {'file_w': file_w, 'file_h': file_h, 'bbox_w': bbox_w, 'bbox_h': bbox_h}
    except Exception as e:
        return {'error': str(e), 'file_w': None, 'file_h': None, 'bbox_w': None, 'bbox_h': None}
    finally:
        os.unlink(tmp_path)

def get_gimp_png_bg_pixel__06570640015ea4861f9ec59a895e97a7_qw35sft2_175b87d7(env, config: dict):
    """Download exported PNG from VM and return background pixel color at top-left corner."""
    import tempfile
    import os
    try:
        from PIL import Image
    except ImportError:
        return {'error': 'PIL not available', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    path = config.get('path', '/home/user/Desktop/verify_blue_v2.png')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path).convert('RGB')
        width, height = img.size
        r, g, b = img.getpixel((10, 10))
        return {'r': int(r), 'g': int(g), 'b': int(b), 'width': width, 'height': height}
    except Exception as e:
        return {'error': str(e), 'r': -1, 'g': -1, 'b': -1, 'width': -1, 'height': -1}
    finally:
        os.unlink(tmp_path)

def get_three_images_brightness__b3992f4af3f3949be116d5f6a9289f79_qw35sft2_c55a41c0(env, config: dict):
    """Get average RGB brightness of three desktop image files from the VM."""
    import tempfile, os
    try:
        from PIL import Image
        import numpy as np
    except ImportError:
        return {'error': 'PIL or numpy not available'}
    paths = {'squirrel': '/home/user/Desktop/squirrel.jpeg', 'panda': '/home/user/Desktop/panda.jpeg', 'heron': '/home/user/Desktop/heron.jpeg'}
    result = {}
    for name, vm_path in paths.items():
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            result[name] = {'avg_brightness': 0.0, 'file_found': False}
            continue
        with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path).convert('RGB')
            arr = np.array(img, dtype=float)
            result[name] = {'avg_brightness': float(arr.mean()), 'file_found': True}
        except Exception as e:
            result[name] = {'avg_brightness': 0.0, 'file_found': False, 'error': str(e)}
        finally:
            os.unlink(tmp_path)
    return result

def get_gimp_jpeg_info__d6fa459b37f87da875d5b6370f3cc58a_qw35sft2_b8f15e4c(env, config: dict):
    """Check if the exported JPEG file exists and return its metadata."""
    from PIL import Image
    file_path = config.get('path', '/home/user/logo.jpg')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'mode': None, 'width': None, 'height': None}
    except Exception:
        return {'exists': False, 'mode': None, 'width': None, 'height': None}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        return {'exists': True, 'mode': img.mode, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'exists': True, 'mode': None, 'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_mode__f54b84b51bb414e1baf7bdf3d2e0e400_qw35sft2_37cb28b8(env, config: dict):
    """Get the PIL mode of an image file from VM.
    Returns dict with 'mode' key ('L' for grayscale, 'RGB' for color, etc.)
    """
    from PIL import Image
    path = config.get('path', '/home/user/Desktop/low_resolution.jpeg')
    dest = config.get('dest', 'mode_check.jpeg')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger_qw35sft2_407382.error('Failed to get file from VM: %s', path)
        return {'mode': 'error', 'error': 'file not found'}
    local_path = os.path.join(env.cache_dir, dest)
    with open(local_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(local_path)
        img.load()
        mode = img.mode
        logger_qw35sft2_407382.info('Image mode: %s, size: %s', mode, img.size)
        return {'mode': mode, 'path': local_path}
    except Exception as e:
        logger_qw35sft2_407382.error('Failed to open image: %s', e)
        return {'mode': 'error', 'error': str(e)}

def get_gimp_xcf_layers__413739fbbe273da47f148920d82ae49b_qw35sft2_3e65ef23(env, config: dict):
    """Parse XCF file from VM and return layer names and count."""
    import struct
    xcf_path = config.get('path', '/home/user/Desktop/white_background.xcf')
    try:
        file_bytes = env.controller.get_file(xcf_path)
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}
    if not file_bytes:
        return {'error': 'File not found or empty', 'layer_names': [], 'layer_count': 0}
    data = bytes(file_bytes)
    if not data.startswith(b'gimp xcf '):
        return {'error': 'Not a valid XCF file', 'layer_names': [], 'layer_count': 0}
    try:
        null_pos = data.index(b'\x00', 9)
        version_str = data[9:null_pos].decode('ascii', errors='replace')
        pos = null_pos + 1
        if pos + 12 > len(data):
            return {'error': 'File too short', 'layer_names': [], 'layer_count': 0}
        pos += 12
        while pos + 8 <= len(data):
            prop_type = struct.unpack('>I', data[pos:pos + 4])[0]
            prop_size = struct.unpack('>I', data[pos + 4:pos + 8])[0]
            pos += 8 + prop_size
            if prop_type == 0:
                break
        try:
            ver_num = int(version_str.lstrip('v') or 0)
        except Exception:
            ver_num = 0
        offset_size = 8 if ver_num >= 11 else 4
        offset_fmt = '>Q' if ver_num >= 11 else '>I'
        layer_offsets = []
        while pos + offset_size <= len(data):
            offset = struct.unpack(offset_fmt, data[pos:pos + offset_size])[0]
            pos += offset_size
            if offset == 0:
                break
            layer_offsets.append(offset)
        layer_names = []
        for off in layer_offsets:
            if off + 16 > len(data):
                continue
            name_off = off + 12
            name_len = struct.unpack('>I', data[name_off:name_off + 4])[0]
            name_off += 4
            if name_len == 0 or name_off + name_len > len(data):
                continue
            name_bytes = data[name_off:name_off + name_len]
            name = name_bytes.rstrip(b'\x00').decode('utf-8', errors='replace')
            layer_names.append(name)
        return {'layer_names': layer_names, 'layer_count': len(layer_names)}
    except Exception as e:
        return {'error': str(e), 'layer_names': [], 'layer_count': 0}

def get_gimp_image_props__857686b0194013b75de5a057cd91e5ad_qw35sft2_8e186571(env, config: dict):
    """Get image mode and dimensions from exported PNG file on the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/scaled_palette_computer.png')
    try:
        from PIL import Image
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            return {'error': 'file not found or empty'}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as f:
            f.write(file_bytes)
            tmp_path = f.name
        try:
            img = Image.open(tmp_path)
            result = {'mode': img.mode, 'width': img.size[0], 'height': img.size[1]}
            img.close()
        finally:
            os.unlink(tmp_path)
        return result
    except Exception as e:
        return {'error': str(e)}

def get_desktop_image_list__1dbc4bff650a401c3959c27fb5e70015_qw35sft2_e03db9b1(env, config: dict):
    """
    Reads ~/Desktop/image_files.txt and runs the reference find command to
    compare against.

    Returns a dict with:
      file_exists     - bool: whether the file exists
      file_lines      - sorted list of non-empty lines from the file
      expected_lines  - sorted list of PNG/JPEG paths found by `find ~`
      all_lines_valid - bool: every line in file ends with .png/.jpg/.jpeg
    """
    file_result = env.controller.run_bash_script('if [ -f ~/Desktop/image_files.txt ]; then  cat ~/Desktop/image_files.txt; else  echo "__NOT_FOUND__"; fi', timeout=10)
    if isinstance(file_result, dict):
        file_stdout = file_result.get('output', file_result.get('stdout', '')).strip()
    else:
        file_stdout = str(file_result).strip()
    if '__NOT_FOUND__' in file_stdout:
        return {'file_exists': False, 'file_lines': [], 'expected_lines': [], 'all_lines_valid': False}
    file_lines = sorted((line.strip() for line in file_stdout.splitlines() if line.strip()))
    image_re = re.compile('\\.(png|jpe?g)$', re.IGNORECASE)
    all_lines_valid = all((image_re.search(line) for line in file_lines)) if file_lines else True
    find_result = env.controller.run_bash_script('find ~ -type f \\( -iname "*.png" -o -iname "*.jpg" -o -iname "*.jpeg" \\) 2>/dev/null', timeout=30)
    if isinstance(find_result, dict):
        find_stdout = find_result.get('output', find_result.get('stdout', '')).strip()
    else:
        find_stdout = str(find_result).strip()
    expected_lines = sorted((line.strip() for line in find_stdout.splitlines() if line.strip()))
    return {'file_exists': True, 'file_lines': file_lines, 'expected_lines': expected_lines, 'all_lines_valid': all_lines_valid}

def get_transition_and_png__dd35705ee09fd488827a6afdd7cbab81_qw35sft2_03943646(env, config: dict):
    """
    Composite getter: checks both that res.png exists on the Desktop and that
    the source pptx has a slide transition applied to slide 1.
    """
    import tempfile, os
    png_path = config.get('png_path', '/home/user/Desktop/res.png')
    bash_result = env.controller.run_bash_script(f'test -f "{png_path}" && echo "exists" || echo "not_found"', timeout=15)
    if isinstance(bash_result, dict):
        bash_out = bash_result.get('output', '') or bash_result.get('stdout', '') or ''
    else:
        bash_out = str(bash_result) if bash_result else ''
    png_exists = 'exists' in bash_out.strip()
    pptx_path = config.get('pptx_path', '/home/user/Desktop/wssf-project-plan-on-a-page.pptx')
    pptx_bytes = env.controller.get_file(pptx_path)
    has_transition = False
    if pptx_bytes:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp.write(pptx_bytes)
                tmp_path = tmp.name
            from pptx import Presentation
            from pptx.oxml.ns import qn
            prs = Presentation(tmp_path)
            slide = prs.slides[0]
            trans_el = slide._element.find(qn('p:transition'))
            if trans_el is not None:
                fade_el = trans_el.find(qn('p:fade'))
                has_transition = fade_el is not None
            else:
                has_transition = False
        except Exception:
            pass
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)
    return {'png_exists': png_exists, 'has_transition': has_transition}

def get_docx_image_and_alignment__b349befcb19ab6c9d751d2833a5ceb8c_qw35sft2_8ed4cb3f(env, config: dict):
    """Download the docx and return inline image count and whether any image paragraph is center-aligned."""
    import tempfile
    import os
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    WP_INLINE = '{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}inline'
    path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        inline_count = len(doc.inline_shapes)
        has_centered_image = False
        for para in doc.paragraphs:
            if para._p.findall('.//' + WP_INLINE):
                if para.alignment == WD_ALIGN_PARAGRAPH.CENTER:
                    has_centered_image = True
                    break
        return {'image_count': inline_count, 'has_centered_image': has_centered_image}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_image_count__ed4f0e4477e03c56fe74b28d9d6a3444_qw35sft2_acb79876(env, config: dict):
    """Download Viewing_Your_Class_Schedule_and_Textbooks.docx and return the count of inline images."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        inline_count = len(doc.inline_shapes)
        return {'image_count': inline_count}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_image_and_fontsize__285ee55215732bb31d8255226d83a755_qw35sft2_abebb093(env, config: dict):
    """Download the docx and return inline image count and font size of the Figure 1 caption paragraph."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        inline_count = len(doc.inline_shapes)
        caption_font_size_pt = None
        for para in doc.paragraphs:
            if para.text.startswith('Figure 1'):
                for run in para.runs:
                    if run.font.size is not None:
                        caption_font_size_pt = round(run.font.size / 12700)
                        break
                break
        return {'image_count': inline_count, 'caption_font_size_pt': caption_font_size_pt}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_image_and_pdf__cba3b2d7a3f4e1f51b9a84ac71e8e772_qw35sft2_e2f152b3(env, config: dict):
    """Return inline image count from docx and whether the exported PDF exists on the desktop."""
    import tempfile
    import os
    from docx import Document
    docx_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    pdf_path = config.get('pdf_path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.pdf')
    pdf_check = env.controller.run_bash_script(f'test -f "{pdf_path}" && echo EXISTS || echo MISSING', timeout=10)
    pdf_exists = False
    if isinstance(pdf_check, dict):
        output = pdf_check.get('output', '') or pdf_check.get('stdout', '')
        pdf_exists = 'EXISTS' in str(output)
    elif isinstance(pdf_check, str):
        pdf_exists = 'EXISTS' in pdf_check
    file_bytes = env.controller.get_file(docx_path)
    image_count = 0
    if file_bytes:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            image_count = len(doc.inline_shapes)
        except Exception:
            pass
        finally:
            os.unlink(tmp_path)
    return {'image_count': image_count, 'pdf_exists': pdf_exists}

def get_docx_image_and_pagenumbers__6390ff8ee787412726544ea6ba43db0d_qw35sft2_fd1dbc8e(env, config: dict):
    """Download the docx and return inline image count and whether footer contains page numbers."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        inline_count = len(doc.inline_shapes)
        W_INSTR = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instrText'
        has_page_numbers = False
        for section in doc.sections:
            footer = section.footer
            for para in footer.paragraphs:
                xml = para._p.xml
                if 'PAGE' in xml:
                    has_page_numbers = True
                    break
            if has_page_numbers:
                break
        return {'image_count': inline_count, 'has_page_numbers': has_page_numbers}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_small_jpeg_size__98edc79a225cc1dd6b5d8a23724241ab_qw35sft2_1b132276(env, config: dict):
    """Get size in bytes of the aggressively compressed JPEG file from the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/compressed.jpeg')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'size_bytes': -1}
    return {'size_bytes': len(file_bytes)}

def get_image_dimensions__b392913b04d99d3e7513b78f2dedf15d_qw35sft2_078d6e64(env, config: dict):
    """Download an image from the VM and return its pixel dimensions."""
    import tempfile, os
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'exists': False}
    suffix = os.path.splitext(path)[1] or '.png'
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        width, height = img.size
        return {'width': width, 'height': height, 'exists': True}
    except Exception as e:
        return {'error': str(e), 'exists': False}
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

def get_image_size_state__8cc8ca1163a74ceeec7ab0fdb85713f3_qw35sft2_5465ba2b(env, config: dict):
    """Get file size and image dimensions of compressed.jpeg from the VM Desktop."""
    import tempfile
    import os
    path = config.get('path', '/home/user/Desktop/compressed.jpeg')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'size_bytes': -1, 'width': -1, 'height': -1}
    with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        from PIL import Image
        img = Image.open(tmp_path)
        return {'size_bytes': len(file_bytes), 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'error': str(e), 'size_bytes': len(file_bytes), 'width': -1, 'height': -1}
    finally:
        os.unlink(tmp_path)

def get_image_dimensions__89f1952738bca9658039f047c420ba2d_qw35sft2_38f42e78(env, config: dict):
    """Get the width and height of an image file from the VM using Python PIL."""
    path = config.get('path', '/home/user/Desktop/pic.jpg')
    script = f"""python3 -c "from PIL import Image; img = Image.open('{path}'); print(img.size[0], img.size[1])" 2>/dev/null || echo ERROR"""
    result = env.controller.run_bash_script(script, timeout=20)
    output = (result.get('output') or '').strip() if isinstance(result, dict) else str(result).strip()
    if 'ERROR' in output or not output:
        return {'error': 'Could not read image dimensions', 'width': None, 'height': None}
    parts = output.split()
    if len(parts) < 2:
        return {'error': 'Unexpected output format', 'width': None, 'height': None}
    return {'width': int(parts[0]), 'height': int(parts[1])}

def get_jpeg_file_size__19b267768d7fcf0b434b3cca9c02b15d_qw35sft2_a7b7be93(env, config: dict):
    """Get size in bytes of the compressed JPEG file from the VM Desktop."""
    path = config.get('path', '/home/user/Desktop/compressed.jpeg')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'size_bytes': -1}
    return {'size_bytes': len(file_bytes)}

def get_jpg_move_state__7427978e92f6fc0e5652a4713261a5a8_qw35sft2_821f0d4a(env, config: dict):
    """Check jpgs in cpjpg and remaining jpgs in photos (move verification)."""
    try:
        result_cpjpg = env.controller.run_bash_script("find /home/user/Desktop/cpjpg -maxdepth 1 -type f -iname '*.jpg' | wc -l", timeout=30)
        result_photos = env.controller.run_bash_script("find /home/user/Desktop/photos -type f -iname '*.jpg' | wc -l", timeout=30)
        cpjpg_count = int(result_cpjpg.get('output', '0').strip()) if result_cpjpg else 0
        photos_jpg_remaining = int(result_photos.get('output', '0').strip()) if result_photos else -1
        return {'cpjpg_count': cpjpg_count, 'photos_jpg_remaining': photos_jpg_remaining}
    except Exception as e:
        return {'error': str(e), 'cpjpg_count': 0, 'photos_jpg_remaining': -1}

def get_jpg_png_copy_state__ec3ddc36152b3d2fb4aee4f74f969e31_qw35sft2_2c06a286(env, config: dict):
    """Check how many .jpg files are in cpjpg and .png files in cppng."""
    try:
        result_jpg = env.controller.run_bash_script("find /home/user/Desktop/cpjpg -maxdepth 1 -type f -iname '*.jpg' | wc -l", timeout=30)
        result_png = env.controller.run_bash_script("find /home/user/Desktop/cppng -maxdepth 1 -type f -iname '*.png' 2>/dev/null | wc -l", timeout=30)
        jpg_count = int(result_jpg.get('output', '0').strip()) if result_jpg else 0
        png_count = int(result_png.get('output', '0').strip()) if result_png else 0
        return {'jpg_count': jpg_count, 'png_count': png_count}
    except Exception as e:
        return {'error': str(e), 'jpg_count': 0, 'png_count': 0}

def get_jpg_copy_with_count__58fb65f54a1152c1f3aecb552e15d932_qw35sft2_28376ab0(env, config: dict):
    """Check jpgs in cpjpg and content of jpg_count.txt."""
    try:
        result_jpg = env.controller.run_bash_script("find /home/user/Desktop/cpjpg -maxdepth 1 -type f -iname '*.jpg' | wc -l", timeout=30)
        jpg_count = int(result_jpg.get('output', '0').strip()) if result_jpg else 0
        result_file = env.controller.run_bash_script("cat /home/user/Desktop/jpg_count.txt 2>/dev/null || echo ''", timeout=30)
        count_file_value = result_file.get('output', '').strip() if result_file else ''
        return {'jpg_count': jpg_count, 'count_file_value': count_file_value}
    except Exception as e:
        return {'error': str(e), 'jpg_count': 0, 'count_file_value': ''}

def get_jpg_filelist__b885c2e3b122c9010091a38454018836_qw35sft2_048d9885(env, config: dict):
    """Check jpgs in cpjpg and content of filelist.txt."""
    try:
        result_jpg = env.controller.run_bash_script("find /home/user/Desktop/cpjpg -maxdepth 1 -type f -iname '*.jpg' | wc -l", timeout=30)
        jpg_count = int(result_jpg.get('output', '0').strip()) if result_jpg else 0
        result_file = env.controller.run_bash_script("cat /home/user/Desktop/filelist.txt 2>/dev/null || echo ''", timeout=30)
        filelist_content = result_file.get('output', '').strip() if result_file else ''
        return {'jpg_count': jpg_count, 'filelist_content': filelist_content}
    except Exception as e:
        return {'error': str(e), 'jpg_count': 0, 'filelist_content': ''}

def get_vacation_jpg_copy__91c6f86e45abe96db716a4e3d1072be2_qw35sft2_fe78b515(env, config: dict):
    """List files in cpjpg and count jpgs from vacation vs all photos."""
    try:
        result_ls = env.controller.run_bash_script("ls /home/user/Desktop/cpjpg/ 2>/dev/null || echo ''", timeout=30)
        result_vacation = env.controller.run_bash_script("find /home/user/Desktop/cpjpg -maxdepth 1 -type f -iname '*.jpg' | wc -l", timeout=30)
        listing = result_ls.get('output', '').strip() if result_ls else ''
        jpg_count = int(result_vacation.get('output', '0').strip()) if result_vacation else 0
        return {'listing': listing, 'jpg_count': jpg_count}
    except Exception as e:
        return {'error': str(e), 'listing': '', 'jpg_count': 0}

def get_vlc_image_adjust__1e9b265786c948c671445fe3130d0b36_qw35sft2_687fc6f9(env, config: dict):
    """Read VLC config and return image-adjust enable state and brightness."""
    config_path = '/home/user/.config/vlc/vlcrc'
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'error': 'vlcrc not found', 'adjust_enabled': False, 'brightness': 1.0}
    content = file_bytes.decode('utf-8', errors='replace')
    adjust_enabled = False
    brightness = 1.0
    for line in content.splitlines():
        line = line.strip()
        if line.startswith('#') or '=' not in line:
            continue
        key, _, val = line.partition('=')
        key = key.strip()
        val = val.strip()
        if key == 'video-filter':
            adjust_enabled = 'adjust' in val.lower()
        elif key in ('brightness', 'adjust-brightness'):
            try:
                brightness = float(val)
            except (ValueError, TypeError):
                pass
    return {'adjust_enabled': adjust_enabled, 'brightness': brightness}

def get_vlc_image_adjust__33b1c5d9e144110cf196db624efc6d81_qw35sft2_573798c0(env, config: dict):
    """Read VLC config and return image-adjust enable state."""
    config_path = '/home/user/.config/vlc/vlcrc'
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'error': 'vlcrc not found', 'adjust_enabled': False}
    content = file_bytes.decode('utf-8', errors='replace')
    adjust_enabled = False
    for line in content.splitlines():
        line = line.strip()
        if line.startswith('#') or '=' not in line:
            continue
        key, _, val = line.partition('=')
        key = key.strip()
        val = val.strip()
        if key == 'video-filter':
            adjust_enabled = 'adjust' in val.lower()
    return {'adjust_enabled': adjust_enabled}

def get_vlc_image_adjust__493da9b8d1a20a26a98a5cc60b751e43_qw35sft2_2106ab38(env, config: dict):
    """Read VLC config and return image-adjust enable state and saturation."""
    config_path = '/home/user/.config/vlc/vlcrc'
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'error': 'vlcrc not found', 'adjust_enabled': False, 'saturation': 1.0}
    content = file_bytes.decode('utf-8', errors='replace')
    adjust_enabled = False
    saturation = 1.0
    for line in content.splitlines():
        line = line.strip()
        if line.startswith('#') or '=' not in line:
            continue
        key, _, val = line.partition('=')
        key = key.strip()
        val = val.strip()
        if key == 'video-filter':
            adjust_enabled = 'adjust' in val.lower()
        elif key in ('saturation', 'adjust-saturation'):
            try:
                saturation = float(val)
            except (ValueError, TypeError):
                pass
    return {'adjust_enabled': adjust_enabled, 'saturation': saturation}

def get_vlc_image_adjust__c451d2a7e7a5b8e510cce0c34da4325d_qw35sft2_3cf3258a(env, config: dict):
    """Read VLC config and return image-adjust and brightness-threshold states."""
    config_path = '/home/user/.config/vlc/vlcrc'
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'error': 'vlcrc not found', 'adjust_enabled': False, 'brightness_threshold': False}
    content = file_bytes.decode('utf-8', errors='replace')
    adjust_enabled = False
    brightness_threshold = False
    for line in content.splitlines():
        line = line.strip()
        if line.startswith('#') or '=' not in line:
            continue
        key, _, val = line.partition('=')
        key = key.strip()
        val = val.strip()
        if key == 'video-filter':
            adjust_enabled = 'adjust' in val.lower()
        elif key in ('brightness-threshold', 'adjust-brightness-threshold'):
            brightness_threshold = val in ('1', 'true', 'True')
    return {'adjust_enabled': adjust_enabled, 'brightness_threshold': brightness_threshold}

def get_vlc_image_adjust__ab50e89f0aef704ab57d4c9d2579417f_qw35sft2_83e6738b(env, config: dict):
    """Read VLC config and return image-adjust enable state and contrast."""
    config_path = '/home/user/.config/vlc/vlcrc'
    file_bytes = env.controller.get_file(config_path)
    if not file_bytes:
        return {'error': 'vlcrc not found', 'adjust_enabled': False, 'contrast': 1.0}
    content = file_bytes.decode('utf-8', errors='replace')
    adjust_enabled = False
    contrast = 1.0
    for line in content.splitlines():
        line = line.strip()
        if line.startswith('#') or '=' not in line:
            continue
        key, _, val = line.partition('=')
        key = key.strip()
        val = val.strip()
        if key == 'video-filter':
            adjust_enabled = 'adjust' in val.lower()
        elif key in ('contrast', 'adjust-contrast'):
            try:
                contrast = float(val)
            except (ValueError, TypeError):
                pass
    return {'adjust_enabled': adjust_enabled, 'contrast': contrast}
