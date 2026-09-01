"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_vlc_snapshot_exists__786ac827ab684b501ddc668a08ec2cfa', 'get_vlc_snapshot_subdir__5366805ab8e670a4853b7e253af708d2', 'get_vlc_snapshot_exists__441247156205cdab6fb75baec6278a39', 'get_audio_extraction_state__e0d33b8eb904792c666a0fe3a33dc4f5', 'get_vlc_snapshot_path__5e0183d902d6d20115487cceb9d1bd09', 'get_vlc_dual_snapshot__1118aad2dd728a496771a4813a51c0c2', 'get_video_dimensions__92a1bbba6e28a320b340946ad5c7a75d', 'get_vlc_snapshot_file__da77e4abbc92d5a7a2f84a9516d6d0b6', 'get_vlc_max_volume__ccc1da9663fce350237c6ebebb163164', 'get_video_dimensions__22c3b2acb75b2231bd0ecc7a3dd0e178', 'get_audio_file_info__7c8156c2a09b3bba70e2cfa0aa1a4dcc', 'get_audio_file_info__9e6e6e9addabc1fa910c5d7bd33b2f51', 'get_vlc_network_cache__72b69002c9e414b44114ece1b158222e', 'get_vlc_multi_settings__be291c8b2f31ab750f3a52d157cbc9fd', 'get_video_nosub_state__bf3c684d8dd921786e84a6a366f46d5a', 'get_audio_file_info__3c3196d8d3d24397883e7bbcfddc8dd0', 'get_video_duration__38add8a045487bc69861547adf883b78_qw35sft2_8f33a5f0', 'get_video_duration__dd6b82d69d19c741223e8dfd4c176567_qw35sft2_dc0b042a', 'get_mp3_exists_mp4_removed__e17edfd562e8913a3a72a3bcb87c9c97_qw35sft2_b3edf76c', 'get_vlc_record_path__33b97504f4caa5a543c607fe8892e403_qw35sft2_f4c58bff', 'get_vlc_max_volume__95b68fd0ad40f238d037990e09029619_qw35sft2_6ffa330a', 'get_vlc_dual_config__de221324409107598f868e822754ff7f_qw35sft2_41a0730f', 'get_vlc_snap_and_rename__25f0d8c371fba607f6a6c0535521c207_qw35sft2_528101f8', 'get_vlc_status_with_random__cd6a6657f06ad537781a5ce037c2b82a_qw35sft2_4c09a37d', 'get_vlc_stream_and_cache__4f78681b35e8e22495902169ef7b5c87_qw35sft2_e2d0457c', 'get_snapshot_and_wallpaper__1e243c15dcedcf4e9a4cacc616a0598b_qw35sft2_1c6f2873', 'get_vlc_recfolder_and_maxvol__517e4427a675aa3407c571ed7524794f_qw35sft2_fe4086dc', 'get_vlc_record_path__e488a41c13481ad2b017a00cb14067ec_qw35sft2_9127e59f', 'get_vlc_max_volume__15cb17bd2e24428fe6d640ea8d5d556a_qw35sft2_f743a5de', 'get_vlc_instance_settings__71c2e6dc1ced259311cab9f07521fd96_qw35sft2_87ec3a01', 'get_mp3_validity__35ef4e9e6e88fc273aff9d3f203be720_qw35sft2_d9b9a999', 'get_vlc_status_with_rate__a87ab91c6345361c3c6e374fc52958bf_qw35sft2_e9f6422b', 'get_vlc_rename_check__bd8f719f06ca22a9a2d291dd693db756_qw35sft2_2dadd661', 'get_vlc_cone_and_oneinstance__6535f71afc4db044467cc2ad41010df9_qw35sft2_6860f345', 'get_wallpaper_and_desktop_snapshot__5688825554c17cb3fb1e24246b369618_qw35sft2_29d38da3', 'get_vlc_stream_and_cone__e0203e1b79d3816d60eca3abd8c45b33_qw35sft2_ce7ce6d3', 'get_vlc_snap_in_captures__6d219cf2312d28ed73dcf8a9fc508b4d_qw35sft2_a2597149', 'get_vlc_fullscreen_and_snapshot__b80790f2f04245c0798c2ba24e3ceaf2_qw35sft2_a2cfc2f7', 'get_vlc_maxvol_and_bgcone__f99d034b7e5202fd562c023da2142c74_qw35sft2_08ba85db', 'get_vlc_playback_prefs__db101724f343bdcc09eb9f11a0e77c94_qw35sft2_2f8d8457', 'get_vlc_effects_dialog__80d9e76aa584c867ce735bfda2f705a1_qw35sft2_e8fc3de4', 'get_vlc_status_with_repeat__dc3aa9f060ebd375fbf717e1b7e05e93_qw35sft2_8bd0dfaf', 'get_mp3_mp4_both_exist__41a079bbfac3c55bed03337726ec0862_qw35sft2_e793a0ca', 'get_vlc_prefs__8b6f202a2cac4c2b2f141e2f0645d358_qw35sft2_f5b01f4a', 'get_vlc_snapshot_in_pictures__4e0abb3a76e2d3729689698c9288f79c_qw35sft2_34ed3771', 'get_vlc_cone_and_volume__78f02ce8f924783e3007881645e6184a_qw35sft2_e8651202', 'get_vlc_stream_and_snapshot__b5aec674ad8f49621ab70bae692967c5_qw35sft2_f27e5749', 'get_vlc_snapshot_dual_loc__b8a50137decabc772595757ced9c456a_qw35sft2_fb36c8cb', 'get_vlc_fullscreen_and_maxvol__ecd5fecf55729059a58a10d83c8c9eeb_qw35sft2_ba75f6f7', 'get_vlc_triple_prefs__29cde5ba2a87fed510792760a2c64d1f_qw35sft2_e351a0e2', 'get_vlc_play_and_record_path__d9c5d9aac51555091a28197b6da34259_qw35sft2_9556c60a', 'get_mp3_file_size__55ceccfa8167b4a354049a9775b28527_qw35sft2_07f0402e', 'get_vlc_maxvol_and_bgcone_expands__d25099f13e7da2513cf33c07be0e7955_qw35sft2_c095702e', 'get_vlc_cone_and_minview__24dc6e88ff00c69857d313c4738a8f31_qw35sft2_e221381a', 'get_vlc_record_path__f613292e2de2a0e0145aaa24fdebbc72_qw35sft2_efd6fb7c', 'get_vlc_stream_and_volume__a51ef1127f9b7bededbd54fd205d555d_qw35sft2_dfc2e8ed', 'get_vlc_snapshot_pictures__cdcbbd90330cd55be87f2b9353b8c43c_qw35sft2_9917fb41', 'get_vlc_fullscreen_and_bgcone__986ee683b9dd7fb42919cda3330ee515_qw35sft2_b9d5fc00', 'get_vlc_loop_prefs__a02aa930219f8e3a3317ea91cceaa2ff_qw35sft2_52cb7a7f', 'get_vlc_status_with_loop__c8546449445cd5e2551a333da05cf053_qw35sft2_dfccd919', 'get_vlc_max_volume__a9382a8b8a5c0cd01bfb6f50454b9a65_qw35sft2_cbb67ad4', 'get_vlc_cone_and_recordpath__914a6c12022e434d9eb3c729da97ace4_qw35sft2_67033980', 'get_mp3_existence__77b822a22aadd202cc6c547979661a88_qw35sft2_ad415aff', 'get_vlc_snap_on_desktop__dc35efec6c78cc3245c9e38e3e5c5d4a_qw35sft2_82878c87', 'get_wallpaper_and_vlc_running__2cdc858f118b65a32b6755071b250267_qw35sft2_22fb9858', 'get_vlc_fullscreen_and_recfolder__27e3d8a3f9225f7e9bc7da9d0f1debd6_qw35sft2_6eeae1b7', 'get_vlc_stream_and_ontop__5d8c2e81cb9645b71051b528d54492c6_qw35sft2_565e3d78', 'get_vlc_prefs__a7d5b4b88d0304ec4fb07fa10130fbc0_qw35sft2_fbd3b038']

def get_vlc_snapshot_exists__786ac827ab684b501ddc668a08ec2cfa(env, config: dict):
    """Check if a VLC snapshot file exists in /home/user/Pictures/."""
    try:
        result = env.controller.run_bash_script('ls /home/user/Pictures/vlcsnap-*.png 2>/dev/null | head -1', timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else ''
        if output and 'vlcsnap-' in output:
            return {'snapshot_exists': True, 'snapshot_path': output}
        return {'snapshot_exists': False, 'snapshot_path': ''}
    except Exception as e:
        logger.error('Error checking VLC snapshot: %s', str(e))
        return {'snapshot_exists': False, 'snapshot_path': '', 'error': str(e)}

def get_vlc_snapshot_subdir__5366805ab8e670a4853b7e253af708d2(env, config: dict):
    """Check if a snapshot file exists in a subdirectory that may need creation."""
    target_path = config.get('path', '/home/user/Pictures/snapshots/interstellar_120s.png')
    target_dir = '/'.join(target_path.rsplit('/', 1)[:-1])
    try:
        dir_result = env.controller.run_bash_script(f"test -d '{target_dir}' && echo 'DIR_EXISTS' || echo 'DIR_MISSING'", timeout=30)
        dir_output = dir_result.get('output', '').strip() if isinstance(dir_result, dict) else str(dir_result).strip()
        dir_exists = 'DIR_EXISTS' in dir_output
        file_result = env.controller.run_bash_script(f"test -f '{target_path}' && file --mime-type -b '{target_path}' && stat -c '%s' '{target_path}'", timeout=30)
        file_output = file_result.get('output', '').strip() if isinstance(file_result, dict) else str(file_result).strip()
        lines = file_output.split('\n')
        if len(lines) >= 2:
            return {'dir_exists': dir_exists, 'file_exists': True, 'mime_type': lines[0].strip(), 'file_size': int(lines[1].strip()), 'path': target_path}
        else:
            return {'dir_exists': dir_exists, 'file_exists': False, 'path': target_path}
    except Exception as e:
        logger.error(f'Error checking snapshot subdir: {e}')
        return {'dir_exists': False, 'file_exists': False, 'path': target_path, 'error': str(e)}

def get_vlc_snapshot_exists__441247156205cdab6fb75baec6278a39(env, config: dict):
    """Check if a VLC snapshot file exists in the specified directory."""
    check_dir = config.get('check_dir', '/home/user/Pictures')
    result = env.controller.run_bash_script(f'ls -1 {check_dir}/vlcsnap-*.png 2>/dev/null | head -5', timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else ''
    if output:
        files = [f.strip() for f in output.split('\n') if f.strip()]
        return {'exists': True, 'count': len(files), 'files': files}
    return {'exists': False, 'count': 0, 'files': []}

def get_audio_extraction_state__e0d33b8eb904792c666a0fe3a33dc4f5(env, config: dict):
    """Check if audio was extracted from video correctly."""
    target_path = config.get('path', '/home/user/audio.mp3')
    check_exists = env.controller.run_bash_script(f"test -f {target_path} && echo 'yes' || echo 'no'", timeout=10)
    file_exists = check_exists.get('output', '').strip() == 'yes'
    if not file_exists:
        return {'file_exists': False, 'is_audio': False, 'has_content': False}
    check_audio = env.controller.run_bash_script(f'ffprobe -v quiet -show_entries stream=codec_type -of csv=p=0 {target_path} 2>/dev/null', timeout=15)
    audio_output = check_audio.get('output', '').strip()
    is_audio = 'audio' in audio_output
    check_size = env.controller.run_bash_script(f'stat -c %s {target_path} 2>/dev/null', timeout=10)
    try:
        file_size = int(check_size.get('output', '0').strip())
        has_content = file_size > 1024
    except ValueError:
        has_content = False
    return {'file_exists': file_exists, 'is_audio': is_audio, 'has_content': has_content}

def get_vlc_snapshot_path__5e0183d902d6d20115487cceb9d1bd09(env, config: dict):
    """Get VLC config file to check snapshot-path setting."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\\\AppData\\\\Roaming\\\\vlc\\\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlcrc'))
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_dual_snapshot__1118aad2dd728a496771a4813a51c0c2(env, config: dict):
    """Check if two snapshot files exist at specified paths."""
    paths = config.get('paths', ['/home/user/Desktop/frame1.png', '/home/user/Desktop/frame2.png'])
    results = {}
    for (i, path) in enumerate(paths):
        key = f'file_{i}'
        try:
            check_cmd = f"test -f '{path}' && file --mime-type -b '{path}' && stat -c '%s' '{path}'"
            result = env.controller.run_bash_script(check_cmd, timeout=30)
            output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
            lines = output.split('\n')
            if len(lines) >= 2:
                results[key] = {'exists': True, 'mime_type': lines[0].strip(), 'file_size': int(lines[1].strip()), 'path': path}
            else:
                results[key] = {'exists': False, 'path': path}
        except Exception as e:
            logger.error(f'Error checking file {path}: {e}')
            results[key] = {'exists': False, 'path': path, 'error': str(e)}
    return results

def get_video_dimensions__92a1bbba6e28a320b340946ad5c7a75d(env, config: dict):
    """Get video file dimensions using ffprobe."""
    path = config.get('path', '')
    result = env.controller.run_bash_script(f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height -of csv=p=0 '{path}' 2>/dev/null", timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    if not output:
        return {'exists': False, 'error': 'File not found or not a valid video'}
    parts = output.split(',')
    if len(parts) >= 2:
        try:
            return {'exists': True, 'width': int(parts[0]), 'height': int(parts[1])}
        except (ValueError, IndexError):
            return {'exists': False, 'error': 'Could not parse video dimensions'}
    return {'exists': False, 'error': 'Unexpected ffprobe output'}

def get_vlc_snapshot_file__da77e4abbc92d5a7a2f84a9516d6d0b6(env, config: dict):
    """Check if a snapshot file exists at the specified path and is a valid PNG image."""
    target_path = config.get('path', '/home/user/Documents/trailer_frame.png')
    try:
        result = env.controller.run_bash_script(f"test -f '{target_path}' && file --mime-type -b '{target_path}' && stat -c '%s' '{target_path}'", timeout=30)
        output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
        lines = output.split('\n')
        if len(lines) >= 2:
            mime_type = lines[0].strip()
            file_size = int(lines[1].strip())
            return {'exists': True, 'mime_type': mime_type, 'file_size': file_size, 'path': target_path}
        else:
            return {'exists': False, 'path': target_path, 'error': 'File not found or unreadable'}
    except Exception as e:
        logger.error(f'Error checking snapshot file: {e}')
        return {'exists': False, 'path': target_path, 'error': str(e)}

def get_vlc_max_volume__ccc1da9663fce350237c6ebebb163164(env, config: dict):
    """Get VLC maximum volume setting from vlcrc config file."""
    result = env.controller.run_bash_script("grep -E '^qt-max-volume=' ~/.config/vlc/vlcrc 2>/dev/null || grep -E '#qt-max-volume=' ~/.config/vlc/vlcrc 2>/dev/null || echo 'not_found'", timeout=30)
    stdout = result.get('output', '').strip() if isinstance(result, dict) else ''
    if 'not_found' in stdout or not stdout:
        return {'error': 'vlcrc not found or setting missing', 'raw': stdout}
    is_commented = stdout.startswith('#')
    line = stdout.lstrip('#').strip()
    parts = line.split('=', 1)
    value_str = parts[1].strip() if len(parts) == 2 else ''
    try:
        value = int(value_str)
    except (ValueError, TypeError):
        value = None
    return {'setting_name': 'qt-max-volume', 'value': value, 'is_commented': is_commented, 'raw_line': stdout}

def get_video_dimensions__22c3b2acb75b2231bd0ecc7a3dd0e178(env, config: dict):
    """Get video file dimensions using ffprobe."""
    path = config.get('path', '')
    result = env.controller.run_bash_script(f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height -of csv=p=0 '{path}' 2>/dev/null", timeout=30)
    output = result.get('output', '').strip() if isinstance(result, dict) else str(result).strip()
    if not output:
        return {'exists': False, 'error': 'File not found or not a valid video'}
    parts = output.split(',')
    if len(parts) >= 2:
        try:
            return {'exists': True, 'width': int(parts[0]), 'height': int(parts[1])}
        except (ValueError, IndexError):
            return {'exists': False, 'error': 'Could not parse video dimensions'}
    return {'exists': False, 'error': 'Unexpected ffprobe output'}

def get_audio_file_info__7c8156c2a09b3bba70e2cfa0aa1a4dcc(env, config: dict):
    """Get audio file information using ffprobe on the VM."""
    path = config.get('path', '')
    check_result = env.controller.run_bash_script(f'test -f "{path}" && echo "EXISTS" || echo "NOT_FOUND"', timeout=10)
    exists_output = check_result.get('output', '').strip() if check_result else ''
    if exists_output != 'EXISTS':
        return {'exists': False}
    ffprobe_cmd = f'ffprobe -v quiet -print_format json -show_format -show_streams "{path}"'
    probe_result = env.controller.run_bash_script(ffprobe_cmd, timeout=30)
    output = probe_result.get('output', '') if probe_result else ''
    try:
        info = json.loads(output)
        audio_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'audio']
        video_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'video']
        fmt = info.get('format', {})
        return {'exists': True, 'format_name': fmt.get('format_name', ''), 'duration': float(fmt.get('duration', 0)), 'size': int(fmt.get('size', 0)), 'has_audio': len(audio_streams) > 0, 'has_video': len(video_streams) > 0, 'audio_codec': audio_streams[0].get('codec_name', '') if audio_streams else '', 'sample_rate': audio_streams[0].get('sample_rate', '') if audio_streams else ''}
    except (json.JSONDecodeError, ValueError, KeyError) as e:
        logger.error(f'Failed to parse ffprobe output: {e}')
        return {'exists': True, 'parse_error': True}

def get_audio_file_info__9e6e6e9addabc1fa910c5d7bd33b2f51(env, config: dict):
    """Get audio file information using ffprobe on the VM."""
    path = config.get('path', '')
    check_result = env.controller.run_bash_script(f'test -f "{path}" && echo "EXISTS" || echo "NOT_FOUND"', timeout=10)
    exists_output = check_result.get('output', '').strip() if check_result else ''
    if exists_output != 'EXISTS':
        return {'exists': False}
    ffprobe_cmd = f'ffprobe -v quiet -print_format json -show_format -show_streams "{path}"'
    probe_result = env.controller.run_bash_script(ffprobe_cmd, timeout=30)
    output = probe_result.get('output', '') if probe_result else ''
    try:
        info = json.loads(output)
        audio_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'audio']
        video_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'video']
        fmt = info.get('format', {})
        return {'exists': True, 'format_name': fmt.get('format_name', ''), 'duration': float(fmt.get('duration', 0)), 'size': int(fmt.get('size', 0)), 'has_audio': len(audio_streams) > 0, 'has_video': len(video_streams) > 0, 'audio_codec': audio_streams[0].get('codec_name', '') if audio_streams else '', 'sample_rate': audio_streams[0].get('sample_rate', '') if audio_streams else ''}
    except (json.JSONDecodeError, ValueError, KeyError) as e:
        logger.error(f'Failed to parse ffprobe output: {e}')
        return {'exists': True, 'parse_error': True}

def get_vlc_network_cache__72b69002c9e414b44114ece1b158222e(env, config: dict):
    """Get VLC config file to check network-caching setting."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\\\AppData\\\\Roaming\\\\vlc\\\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlcrc'))
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_multi_settings__be291c8b2f31ab750f3a52d157cbc9fd(env, config: dict):
    """Get multiple VLC Qt settings from the vlcrc config file."""
    dest = config.get('dest', 'vlcrc')
    platform = env.vm_platform if hasattr(env, 'vm_platform') else 'linux'
    if platform == 'linux':
        vlc_config_path = '/home/user/.config/vlc/vlcrc'
    elif platform == 'darwin':
        vlc_config_path = '/Users/user/Library/Preferences/org.videolan.vlc/vlcrc'
    else:
        vlc_config_path = 'C:\\Users\\user\\AppData\\Roaming\\vlc\\vlcrc'
    file_bytes = env.controller.get_file(vlc_config_path)
    if not file_bytes:
        return {'error': 'VLC config file not found'}
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    config_text = file_bytes.decode('utf-8', errors='replace')
    result = {}
    qt_max_volume = '125'
    qt_bgcone = '1'
    for line in config_text.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            qt_max_volume = line.split('=')[-1].strip()
        if 'qt-bgcone=' in line:
            qt_bgcone = line.split('=')[-1].strip()
    result['qt_max_volume'] = qt_max_volume
    result['qt_bgcone'] = qt_bgcone
    return result

def get_video_nosub_state__bf3c684d8dd921786e84a6a366f46d5a(env, config: dict):
    """Check if video was created without subtitle streams."""
    target_path = config.get('path', '/home/user/video_clean.mp4')
    check_exists = env.controller.run_bash_script(f"test -f {target_path} && echo 'yes' || echo 'no'", timeout=10)
    file_exists = check_exists.get('output', '').strip() == 'yes'
    if not file_exists:
        return {'file_exists': False, 'has_video': False, 'has_audio': False, 'no_subtitles': False}
    check_streams = env.controller.run_bash_script(f'ffprobe -v quiet -show_entries stream=codec_type -of csv=p=0 {target_path} 2>/dev/null', timeout=15)
    streams_output = check_streams.get('output', '').strip()
    stream_types = [s.strip() for s in streams_output.split('\n') if s.strip()]
    has_video = 'video' in stream_types
    has_audio = 'audio' in stream_types
    no_subtitles = 'subtitle' not in stream_types
    return {'file_exists': file_exists, 'has_video': has_video, 'has_audio': has_audio, 'no_subtitles': no_subtitles}

def get_audio_file_info__3c3196d8d3d24397883e7bbcfddc8dd0(env, config: dict):
    """Get audio file information using ffprobe on the VM."""
    path = config.get('path', '')
    check_result = env.controller.run_bash_script(f'test -f "{path}" && echo "EXISTS" || echo "NOT_FOUND"', timeout=10)
    exists_output = check_result.get('output', '').strip() if check_result else ''
    if exists_output != 'EXISTS':
        return {'exists': False}
    ffprobe_cmd = f'ffprobe -v quiet -print_format json -show_format -show_streams "{path}"'
    probe_result = env.controller.run_bash_script(ffprobe_cmd, timeout=30)
    output = probe_result.get('output', '') if probe_result else ''
    try:
        info = json.loads(output)
        audio_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'audio']
        video_streams = [s for s in info.get('streams', []) if s.get('codec_type') == 'video']
        fmt = info.get('format', {})
        return {'exists': True, 'format_name': fmt.get('format_name', ''), 'duration': float(fmt.get('duration', 0)), 'size': int(fmt.get('size', 0)), 'has_audio': len(audio_streams) > 0, 'has_video': len(video_streams) > 0, 'audio_codec': audio_streams[0].get('codec_name', '') if audio_streams else '', 'sample_rate': audio_streams[0].get('sample_rate', '') if audio_streams else ''}
    except (json.JSONDecodeError, ValueError, KeyError) as e:
        logger.error(f'Failed to parse ffprobe output: {e}')
        return {'exists': True, 'parse_error': True}

def get_video_duration__38add8a045487bc69861547adf883b78_qw35sft2_8f33a5f0(env, config: dict):
    """Check if a clipped video exists and return its duration via ffprobe."""
    path = config.get('path', '/home/user/clip_3s.mp4')
    exists_result = env.controller.run_bash_script(f'test -f "{path}" && echo EXISTS || echo MISSING', timeout=10)
    if not exists_result or 'EXISTS' not in exists_result.get('output', ''):
        return {'duration': -1.0, 'exists': False}
    dur_result = env.controller.run_bash_script(f'ffprobe -v quiet -show_entries format=duration -of csv=p=0 "{path}" 2>&1', timeout=30)
    output = dur_result.get('output', '').strip() if dur_result else ''
    try:
        duration = float(output)
        return {'duration': duration, 'exists': True}
    except (ValueError, TypeError):
        return {'duration': -1.0, 'exists': True}

def get_video_duration__dd6b82d69d19c741223e8dfd4c176567_qw35sft2_dc0b042a(env, config: dict):
    """Check if a trimmed video exists and return its duration via ffprobe."""
    path = config.get('path', '/home/user/trimmed.mp4')
    exists_result = env.controller.run_bash_script(f'test -f "{path}" && echo EXISTS || echo MISSING', timeout=10)
    if not exists_result or 'EXISTS' not in exists_result.get('output', ''):
        return {'duration': -1.0, 'exists': False}
    dur_result = env.controller.run_bash_script(f'ffprobe -v quiet -show_entries format=duration -of csv=p=0 "{path}" 2>&1', timeout=30)
    output = dur_result.get('output', '').strip() if dur_result else ''
    try:
        duration = float(output)
        return {'duration': duration, 'exists': True}
    except (ValueError, TypeError):
        return {'duration': -1.0, 'exists': True}

def get_mp3_exists_mp4_removed__e17edfd562e8913a3a72a3bcb87c9c97_qw35sft2_b3edf76c(env, config: dict):
    """Check if MP3 was created and original MP4 was deleted."""
    mp3_bytes = env.controller.get_file(MP3_PATH_qw35sft2_fd9ea1)
    mp4_bytes = env.controller.get_file(MP4_PATH_qw35sft2_fd9ea1)
    return {'mp3_exists': mp3_bytes is not None and len(mp3_bytes) > 0, 'mp4_exists': mp4_bytes is not None and len(mp4_bytes) > 0}

def get_vlc_record_path__33b97504f4caa5a543c607fe8892e403_qw35sft2_f4c58bff(env, config: dict):
    """Get VLC recording path from vlcrc configuration file."""
    try:
        result = env.controller.run_bash_script("grep -E '^record-path=' /home/user/.config/vlc/vlcrc 2>/dev/null | head -1", timeout=10)
        if result is None:
            return {'record_path': None, 'error': 'No result from bash script'}
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        output = output.strip()
        if output and '=' in output:
            path = output.split('=', 1)[1].strip()
            return {'record_path': path}
        return {'record_path': None}
    except Exception as e:
        return {'record_path': None, 'error': str(e)}

def get_vlc_max_volume__95b68fd0ad40f238d037990e09029619_qw35sft2_6ffa330a(env, config: dict):
    """Read qt-max-volume from VLC's vlcrc config file."""
    result = {'qt_max_volume': '125'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_05799f.warning('vlcrc not found or empty')
        return result
    for line in vlcrc_bytes.decode('utf-8', errors='ignore').split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
            break
    logger_qw35sft2_05799f.info('VLC qt-max-volume: %s', result)
    return result

def get_vlc_dual_config__de221324409107598f868e822754ff7f_qw35sft2_41a0730f(env, config: dict):
    """
    Reads VLC configuration file and returns qt-minimal-view and qt-max-volume values.
    Used for two-objective verification: minimal view enabled AND max volume set.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    else:
        return {'error': f'Unsupported OS: {os_type}'}
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'vlcrc not found'}
    config_text = content.decode('utf-8', errors='replace')
    qt_minimal_view = '0'
    qt_max_volume = '125'
    for line in config_text.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-minimal-view=' in line:
            qt_minimal_view = line.split('=')[-1].strip()
        if 'qt-max-volume=' in line:
            qt_max_volume = line.split('=')[-1].strip()
    return {'qt_minimal_view': qt_minimal_view, 'qt_max_volume': qt_max_volume}

def get_vlc_snap_and_rename__25f0d8c371fba607f6a6c0535521c207_qw35sft2_528101f8(env, config: dict):
    """Check if interstellar.png on Desktop and video renamed to trailer.mp4."""
    snap_bytes = env.controller.get_file('/home/user/Desktop/interstellar.png')
    snap_exists = bool(snap_bytes and len(snap_bytes) > 1000)
    trailer_result = env.controller.run_bash_script('test -f "/home/user/Desktop/trailer.mp4" && echo "exists" || echo "missing"', timeout=10)
    if isinstance(trailer_result, dict):
        trailer_output = trailer_result.get('output', trailer_result.get('stdout', ''))
    else:
        trailer_output = str(trailer_result)
    trailer_renamed = 'exists' in trailer_output
    orig_result = env.controller.run_bash_script('test -f "/home/user/Desktop/Interstellar Movie - Official Trailer.mp4" && echo "exists" || echo "missing"', timeout=10)
    if isinstance(orig_result, dict):
        orig_output = orig_result.get('output', orig_result.get('stdout', ''))
    else:
        orig_output = str(orig_result)
    original_gone = 'missing' in orig_output
    return {'snapshot_on_desktop': snap_exists, 'trailer_renamed': trailer_renamed, 'original_removed': original_gone}

def get_vlc_status_with_random__cd6a6657f06ad537781a5ce037c2b82a_qw35sft2_4c09a37d(env, config: dict):
    """Get VLC playing status and random/shuffle mode from the HTTP status API."""
    import requests
    from xml.etree import ElementTree
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
    except Exception as e:
        return {'error': str(e)}
    if response.status_code != 200:
        return {'error': f'HTTP {response.status_code}'}
    try:
        tree = ElementTree.fromstring(response.content)
    except Exception as e:
        return {'error': f'XML parse error: {e}'}
    state = tree.findtext('state', default='stopped')
    random = tree.findtext('random', default='false')
    file_name = None
    file_xpath_list = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
    for xp in file_xpath_list:
        elem = tree.find(xp)
        if elem is not None and elem.text:
            file_name = os.path.basename(elem.text)
            break
    return {'state': state, 'random': random, 'file_name': file_name}

def get_vlc_stream_and_cache__4f78681b35e8e22495902169ef7b5c87_qw35sft2_e2d0457c(env, config: dict):
    """Read VLC recent MRL list and network-caching setting from vlcrc."""
    mrl_content = ''
    try:
        mrl_bytes = env.controller.get_file('/home/user/.config/vlc/vlc-qt-interface.conf')
        if mrl_bytes:
            mrl_content = mrl_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_6b91bf.warning('Could not read vlc-qt-interface.conf: %s', e)
    recent_mrl = ''
    for line in mrl_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('list='):
            recent_mrl = stripped[5:].strip()
            break
    vlcrc_content = ''
    try:
        vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
        if vlcrc_bytes:
            vlcrc_content = vlcrc_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_6b91bf.warning('Could not read vlcrc: %s', e)
    network_caching = '1000'
    for line in vlcrc_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('network-caching='):
            network_caching = stripped.split('=', 1)[1].strip()
            break
    return {'recent_mrl': recent_mrl, 'network-caching': network_caching}

def get_snapshot_and_wallpaper__1e243c15dcedcf4e9a4cacc616a0598b_qw35sft2_1c6f2873(env, config: dict):
    """Get both the VLC snapshot count in Pictures and the current wallpaper URI."""
    snap_cmd = 'ls /home/user/Pictures/vlcsnap-*.png 2>/dev/null | wc -l'
    snap_result = env.controller.run_bash_script(snap_cmd, timeout=15)
    if isinstance(snap_result, dict):
        snap_output = snap_result.get('output', snap_result.get('stdout', '0'))
    else:
        snap_output = str(snap_result)
    try:
        snapshot_count = int(snap_output.strip())
    except ValueError:
        snapshot_count = 0
    wall_cmd = 'gsettings get org.gnome.desktop.background picture-uri'
    wall_result = env.controller.run_bash_script(wall_cmd, timeout=15)
    if isinstance(wall_result, dict):
        wall_output = wall_result.get('output', wall_result.get('stdout', ''))
    else:
        wall_output = str(wall_result)
    wallpaper_uri = wall_output.strip().strip('\'"')
    return {'snapshot_count': snapshot_count, 'wallpaper_uri': wallpaper_uri}

def get_vlc_recfolder_and_maxvol__517e4427a675aa3407c571ed7524794f_qw35sft2_fe4086dc(env, config: dict):
    """Get both the recording folder and qt-max-volume settings from VLC's vlcrc config."""
    result = {'input_record_path': None, 'qt_max_volume': '125'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_a330eb.warning('vlcrc not found or empty')
        return result
    vlcrc_content = vlcrc_bytes.decode('utf-8', errors='ignore')
    for line in vlcrc_content.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'input-record-path=' in line:
            result['input_record_path'] = line.split('=', 1)[-1].strip()
        elif 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
    logger_qw35sft2_a330eb.info('VLC recfolder+maxvol state: %s', result)
    return result

def get_vlc_record_path__e488a41c13481ad2b017a00cb14067ec_qw35sft2_9127e59f(env, config: dict):
    """Get VLC recording path from vlcrc configuration file."""
    try:
        result = env.controller.run_bash_script("grep -E '^record-path=' /home/user/.config/vlc/vlcrc 2>/dev/null | head -1", timeout=10)
        if result is None:
            return {'record_path': None, 'error': 'No result from bash script'}
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        output = output.strip()
        if output and '=' in output:
            path = output.split('=', 1)[1].strip()
            return {'record_path': path}
        return {'record_path': None}
    except Exception as e:
        return {'record_path': None, 'error': str(e)}

def get_vlc_max_volume__15cb17bd2e24428fe6d640ea8d5d556a_qw35sft2_f743a5de(env, config: dict):
    """Read qt-max-volume from VLC's vlcrc config file."""
    result = {'qt_max_volume': '125'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_78b50b.warning('vlcrc not found or empty')
        return result
    for line in vlcrc_bytes.decode('utf-8', errors='ignore').split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
            break
    logger_qw35sft2_78b50b.info('VLC qt-max-volume: %s', result)
    return result

def get_vlc_instance_settings__71c2e6dc1ced259311cab9f07521fd96_qw35sft2_87ec3a01(env, config: dict):
    """Read both VLC single-instance settings from vlcrc."""
    result = env.controller.run_bash_script("grep -E '^#?(one-instance|one-instance-when-started-from-file)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
    if isinstance(result, dict):
        text = result.get('output', result.get('stdout', ''))
    else:
        text = str(result) if result else ''
    uncommented = {}
    commented = {}
    for line in text.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        is_comment = line.startswith('#')
        raw = line[1:].strip() if is_comment else line
        if '=' in raw:
            key, _, val = raw.partition('=')
            key = key.strip()
            val = val.strip()
            if is_comment:
                commented[key] = val
            else:
                uncommented[key] = val
    merged = {**commented, **uncommented}
    return {'one_instance': merged.get('one-instance', 'not_set'), 'one_instance_from_file': merged.get('one-instance-when-started-from-file', 'not_set')}

def get_mp3_validity__35ef4e9e6e88fc273aff9d3f203be720_qw35sft2_d9b9a999(env, config: dict):
    """Fetch Baby Justin Bieber.mp3 and check for valid MP3 magic bytes."""
    file_bytes = env.controller.get_file(MP3_PATH_qw35sft2_d3002b)
    if file_bytes is None or len(file_bytes) < 4:
        return {'valid': False, 'size': 0}
    has_id3 = file_bytes[:3] == b'ID3'
    has_sync = file_bytes[0] == 255 and file_bytes[1] & 224 == 224
    return {'valid': has_id3 or has_sync, 'size': len(file_bytes)}

def get_vlc_status_with_rate__a87ab91c6345361c3c6e374fc52958bf_qw35sft2_e9f6422b(env, config: dict):
    """Get VLC playing status and current playback rate from the HTTP status API."""
    import requests
    from xml.etree import ElementTree
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
    except Exception as e:
        return {'error': str(e)}
    if response.status_code != 200:
        return {'error': f'HTTP {response.status_code}'}
    try:
        tree = ElementTree.fromstring(response.content)
    except Exception as e:
        return {'error': f'XML parse error: {e}'}
    state = tree.findtext('state', default='stopped')
    rate = tree.findtext('rate', default='1')
    file_name = None
    file_xpath_list = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
    for xp in file_xpath_list:
        elem = tree.find(xp)
        if elem is not None and elem.text:
            file_name = os.path.basename(elem.text)
            break
    return {'state': state, 'rate': rate, 'file_name': file_name}

def get_vlc_rename_check__bd8f719f06ca22a9a2d291dd693db756_qw35sft2_2dadd661(env, config: dict):
    """
    Check whether the video file was renamed on the Desktop:
    - new file '1984_Apple_Macintosh_Commercial.mp4' exists
    - old file 'flipped_1984_Apple_Macintosh_Commercial.mp4' is gone
    Returns a status string like 'new_exists=1 old_removed=1'.
    """
    vm_ip = env.vm_ip
    port = env.server_port
    cmd = 'python3 -c "import os; n = os.path.exists(\'/home/user/Desktop/1984_Apple_Macintosh_Commercial.mp4\'); o = not os.path.exists(\'/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4\'); print(\'new_exists=%d old_removed=%d\' % (int(n), int(o)))"'
    try:
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': cmd, 'shell': True}, timeout=15)
        if response.status_code == 200:
            return response.json().get('output', '').strip()
        logger_qw35sft2_a85c07.error('Execute returned status %d', response.status_code)
    except Exception as e:
        logger_qw35sft2_a85c07.error('get_vlc_rename_check error: %s', e)
    return 'error'

def get_vlc_cone_and_oneinstance__6535f71afc4db044467cc2ad41010df9_qw35sft2_6860f345(env, config: dict):
    """Read VLC config and return qt-bgcone and one-instance-when-started-from-file settings."""
    try:
        result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")
        config_path = result.get('output', '').strip()
    except Exception:
        config_path = '/home/user/.config/vlc/vlcrc'
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'vlcrc not found'}
    config_text = content.decode('utf-8', errors='replace')
    qt_bgcone = '1'
    one_instance = '1'
    for line in config_text.split('\n'):
        if 'qt-bgcone=' in line:
            qt_bgcone = line.split('=')[-1].strip()
        if 'one-instance-when-started-from-file=' in line:
            one_instance = line.split('=')[-1].strip()
    return {'qt-bgcone': qt_bgcone, 'one-instance-when-started-from-file': one_instance}

def get_wallpaper_and_desktop_snapshot__5688825554c17cb3fb1e24246b369618_qw35sft2_29d38da3(env, config: dict):
    """Get wallpaper URI and check if a VLC snapshot exists on the Desktop."""
    wall_cmd = 'gsettings get org.gnome.desktop.background picture-uri'
    wall_result = env.controller.run_bash_script(wall_cmd, timeout=15)
    if isinstance(wall_result, dict):
        wall_output = wall_result.get('output', wall_result.get('stdout', ''))
    else:
        wall_output = str(wall_result)
    wallpaper_uri = wall_output.strip().strip('\'"')
    desk_cmd = 'ls /home/user/Desktop/vlcsnap-*.png 2>/dev/null | wc -l'
    desk_result = env.controller.run_bash_script(desk_cmd, timeout=15)
    if isinstance(desk_result, dict):
        desk_output = desk_result.get('output', desk_result.get('stdout', '0'))
    else:
        desk_output = str(desk_result)
    try:
        desktop_count = int(desk_output.strip())
    except ValueError:
        desktop_count = 0
    return {'wallpaper_uri': wallpaper_uri, 'desktop_snapshot_count': desktop_count}

def get_vlc_stream_and_cone__e0203e1b79d3816d60eca3abd8c45b33_qw35sft2_ce7ce6d3(env, config: dict):
    """Read VLC recent MRL list and qt-bgcone setting from vlcrc."""
    mrl_content = ''
    try:
        mrl_bytes = env.controller.get_file('/home/user/.config/vlc/vlc-qt-interface.conf')
        if mrl_bytes:
            mrl_content = mrl_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_bf2b87.warning('Could not read vlc-qt-interface.conf: %s', e)
    recent_mrl = ''
    for line in mrl_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('list='):
            recent_mrl = stripped[5:].strip()
            break
    vlcrc_content = ''
    try:
        vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
        if vlcrc_bytes:
            vlcrc_content = vlcrc_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_bf2b87.warning('Could not read vlcrc: %s', e)
    qt_bgcone = '1'
    for line in vlcrc_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('qt-bgcone='):
            qt_bgcone = stripped.split('=', 1)[1].strip()
            break
    return {'recent_mrl': recent_mrl, 'qt-bgcone': qt_bgcone}

def get_vlc_snap_in_captures__6d219cf2312d28ed73dcf8a9fc508b4d_qw35sft2_a2597149(env, config: dict):
    """Check if captures/ folder exists on Desktop and interstellar.png is inside it."""
    folder_result = env.controller.run_bash_script('test -d /home/user/Desktop/captures && echo "dir_exists" || echo "missing"', timeout=10)
    if isinstance(folder_result, dict):
        folder_output = folder_result.get('output', folder_result.get('stdout', ''))
    else:
        folder_output = str(folder_result)
    folder_exists = 'dir_exists' in folder_output
    file_bytes = env.controller.get_file('/home/user/Desktop/captures/interstellar.png')
    file_exists = bool(file_bytes and len(file_bytes) > 1000)
    return {'captures_folder_exists': folder_exists, 'interstellar_in_captures': file_exists}

def get_vlc_fullscreen_and_snapshot__b80790f2f04245c0798c2ba24e3ceaf2_qw35sft2_a2cfc2f7(env, config: dict):
    """Get VLC fullscreen state and whether a snapshot PNG was created in ~/Pictures."""
    result = {}
    screen_size = env.controller.get_vm_screen_size()
    window_size = env.controller.get_vm_window_size(app_class_name='vlc')
    if screen_size and window_size:
        result['is_fullscreen'] = window_size.get('width') == screen_size.get('width') and window_size.get('height') == screen_size.get('height')
    else:
        result['is_fullscreen'] = False
    bash_output = env.controller.run_bash_script('find ~/Pictures -maxdepth 1 -name "*.png" 2>/dev/null | wc -l', timeout=10)
    if bash_output and isinstance(bash_output, dict):
        output_str = bash_output.get('output', '0').strip()
    else:
        output_str = '0'
    try:
        result['snapshot_count'] = int(output_str)
    except (ValueError, TypeError):
        result['snapshot_count'] = 0
    logger_qw35sft2_a3ebf5.info('VLC fullscreen+snapshot state: %s', result)
    return result

def get_vlc_maxvol_and_bgcone__f99d034b7e5202fd562c023da2142c74_qw35sft2_08ba85db(env, config: dict):
    """Read qt-max-volume and qt-bgcone from VLC's vlcrc config file."""
    result = {'qt_max_volume': '125', 'qt_bgcone': '1'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_0975e4.warning('vlcrc not found or empty')
        return result
    for line in vlcrc_bytes.decode('utf-8', errors='ignore').split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
        elif 'qt-bgcone=' in line:
            result['qt_bgcone'] = line.split('=', 1)[-1].strip()
    logger_qw35sft2_0975e4.info('VLC maxvol+bgcone state: %s', result)
    return result

def get_vlc_playback_prefs__db101724f343bdcc09eb9f11a0e77c94_qw35sft2_2f8d8457(env, config: dict):
    """Read VLC instance-from-file and continue-playback settings from vlcrc."""
    result = env.controller.run_bash_script("grep -E '^#?(one-instance-when-started-from-file|continue-playback)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
    if isinstance(result, dict):
        text = result.get('output', result.get('stdout', ''))
    else:
        text = str(result) if result else ''
    uncommented = {}
    commented = {}
    for line in text.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        is_comment = line.startswith('#')
        raw = line[1:].strip() if is_comment else line
        if '=' in raw:
            key, _, val = raw.partition('=')
            key = key.strip()
            val = val.strip()
            if is_comment:
                commented[key] = val
            else:
                uncommented[key] = val
    merged = {**commented, **uncommented}
    return {'one_instance_from_file': merged.get('one-instance-when-started-from-file', 'not_set'), 'continue_playback': merged.get('continue-playback', 'not_set')}

def get_vlc_effects_dialog__80d9e76aa584c867ce735bfda2f705a1_qw35sft2_e8fc3de4(env, config: dict):
    """
    Retrieve the current desktop accessibility tree as a string.
    Used to verify that VLC's 'Adjustments and Effects' dialog is open.
    """
    try:
        tree = env.controller.get_accessibility_tree()
        return tree if tree else ''
    except Exception as e:
        logger_qw35sft2_bfc3c0.error('get_vlc_effects_dialog error: %s', e)
        return ''

def get_vlc_status_with_repeat__dc3aa9f060ebd375fbf717e1b7e05e93_qw35sft2_8bd0dfaf(env, config: dict):
    """Get VLC playing status and repeat mode from the HTTP status API."""
    import requests
    from xml.etree import ElementTree
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
    except Exception as e:
        return {'error': str(e)}
    if response.status_code != 200:
        return {'error': f'HTTP {response.status_code}'}
    try:
        tree = ElementTree.fromstring(response.content)
    except Exception as e:
        return {'error': f'XML parse error: {e}'}
    state = tree.findtext('state', default='stopped')
    repeat = tree.findtext('repeat', default='false')
    file_name = None
    file_xpath_list = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
    for xp in file_xpath_list:
        elem = tree.find(xp)
        if elem is not None and elem.text:
            file_name = os.path.basename(elem.text)
            break
    return {'state': state, 'repeat': repeat, 'file_name': file_name}

def get_mp3_mp4_both_exist__41a079bbfac3c55bed03337726ec0862_qw35sft2_e793a0ca(env, config: dict):
    """Check if both the new MP3 and the original MP4 exist on the desktop."""
    mp3_bytes = env.controller.get_file(MP3_PATH_qw35sft2_4fe03f)
    mp4_bytes = env.controller.get_file(MP4_PATH_qw35sft2_4fe03f)
    return {'mp3_exists': mp3_bytes is not None and len(mp3_bytes) > 0, 'mp4_exists': mp4_bytes is not None and len(mp4_bytes) > 0}

def get_vlc_prefs__8b6f202a2cac4c2b2f141e2f0645d358_qw35sft2_f5b01f4a(env, config: dict):
    """Get VLC recording path and video-title-show preference from vlcrc."""
    try:
        result = env.controller.run_bash_script("grep -E '^(record-path|video-title-show)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
        if result is None:
            return {'record_path': None, 'video_title_show': None, 'error': 'No result from bash script'}
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        record_path = None
        video_title_show = None
        for line in output.strip().splitlines():
            line = line.strip()
            if line.startswith('record-path='):
                record_path = line.split('=', 1)[1].strip()
            elif line.startswith('video-title-show='):
                video_title_show = line.split('=', 1)[1].strip()
        return {'record_path': record_path, 'video_title_show': video_title_show}
    except Exception as e:
        return {'record_path': None, 'video_title_show': None, 'error': str(e)}

def get_vlc_snapshot_in_pictures__4e0abb3a76e2d3729689698c9288f79c_qw35sft2_34ed3771(env, config: dict):
    """Check if a VLC snapshot PNG file exists in the Pictures folder."""
    cmd = 'ls /home/user/Pictures/vlcsnap-*.png 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(cmd, timeout=15)
    if isinstance(result, dict):
        output = result.get('output', result.get('stdout', '0'))
    else:
        output = str(result)
    try:
        count = int(output.strip())
    except ValueError:
        count = 0
    return {'snapshot_count': count}

def get_vlc_cone_and_volume__78f02ce8f924783e3007881645e6184a_qw35sft2_e8651202(env, config: dict):
    """Read VLC config and return qt-bgcone and qt-max-volume settings."""
    try:
        result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")
        config_path = result.get('output', '').strip()
    except Exception:
        config_path = '/home/user/.config/vlc/vlcrc'
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'vlcrc not found'}
    config_text = content.decode('utf-8', errors='replace')
    qt_bgcone = '1'
    qt_max_volume = '125'
    for line in config_text.split('\n'):
        if 'qt-bgcone=' in line:
            qt_bgcone = line.split('=')[-1].strip()
        if 'qt-max-volume=' in line:
            qt_max_volume = line.split('=')[-1].strip()
    return {'qt-bgcone': qt_bgcone, 'qt-max-volume': qt_max_volume}

def get_vlc_stream_and_snapshot__b5aec674ad8f49621ab70bae692967c5_qw35sft2_f27e5749(env, config: dict):
    """Read VLC recent MRL list and check for snapshot files in ~/Pictures and ~/."""
    mrl_content = ''
    try:
        mrl_bytes = env.controller.get_file('/home/user/.config/vlc/vlc-qt-interface.conf')
        if mrl_bytes:
            mrl_content = mrl_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_98c532.warning('Could not read vlc-qt-interface.conf: %s', e)
    recent_mrl = ''
    for line in mrl_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('list='):
            recent_mrl = stripped[5:].strip()
            break
    snapshot_files = []
    try:
        result = env.controller.run_bash_script('find /home/user -maxdepth 3 -name "vlcsnap*.png" 2>/dev/null | head -10', timeout=15)
        output = result.get('output', '').strip() if result else ''
        if output:
            snapshot_files = [f for f in output.split('\n') if f.strip()]
    except Exception as e:
        logger_qw35sft2_98c532.warning('Could not check for snapshot files: %s', e)
    return {'recent_mrl': recent_mrl, 'snapshot_files': snapshot_files, 'snapshot_count': len(snapshot_files)}

def get_vlc_snapshot_dual_loc__b8a50137decabc772595757ced9c456a_qw35sft2_fb36c8cb(env, config: dict):
    """Check if interstellar.png exists on Desktop and in Pictures folder."""
    desktop_bytes = env.controller.get_file('/home/user/Desktop/interstellar.png')
    pictures_bytes = env.controller.get_file('/home/user/Pictures/interstellar.png')
    desktop_ok = bool(desktop_bytes and len(desktop_bytes) > 1000)
    pictures_ok = bool(pictures_bytes and len(pictures_bytes) > 1000)
    return {'desktop_has_interstellar': desktop_ok, 'pictures_has_interstellar': pictures_ok}

def get_vlc_fullscreen_and_maxvol__ecd5fecf55729059a58a10d83c8c9eeb_qw35sft2_ba75f6f7(env, config: dict):
    """Get VLC fullscreen state and the qt-max-volume setting from vlcrc."""
    result = {}
    screen_size = env.controller.get_vm_screen_size()
    window_size = env.controller.get_vm_window_size(app_class_name='vlc')
    if screen_size and window_size:
        result['is_fullscreen'] = window_size.get('width') == screen_size.get('width') and window_size.get('height') == screen_size.get('height')
    else:
        result['is_fullscreen'] = False
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    result['qt_max_volume'] = '125'
    if vlcrc_bytes:
        vlcrc_content = vlcrc_bytes.decode('utf-8', errors='ignore')
        for line in vlcrc_content.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'qt-max-volume=' in line:
                result['qt_max_volume'] = line.split('=', 1)[-1].strip()
                break
    logger_qw35sft2_8bd001.info('VLC fullscreen+maxvol state: %s', result)
    return result

def get_vlc_triple_prefs__29cde5ba2a87fed510792760a2c64d1f_qw35sft2_e351a0e2(env, config: dict):
    """Read VLC one-instance, one-instance-from-file, and continue-playback from vlcrc."""
    result = env.controller.run_bash_script("grep -E '^#?(one-instance|one-instance-when-started-from-file|continue-playback)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
    if isinstance(result, dict):
        text = result.get('output', result.get('stdout', ''))
    else:
        text = str(result) if result else ''
    uncommented = {}
    commented = {}
    for line in text.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        is_comment = line.startswith('#')
        raw = line[1:].strip() if is_comment else line
        if '=' in raw:
            key, _, val = raw.partition('=')
            key = key.strip()
            val = val.strip()
            if is_comment:
                commented[key] = val
            else:
                uncommented[key] = val
    merged = {**commented, **uncommented}
    return {'one_instance': merged.get('one-instance', 'not_set'), 'one_instance_from_file': merged.get('one-instance-when-started-from-file', 'not_set'), 'continue_playback': merged.get('continue-playback', 'not_set')}

def get_vlc_play_and_record_path__d9c5d9aac51555091a28197b6da34259_qw35sft2_9556c60a(env, config: dict):
    """Get VLC playing status and vlcrc recording directory setting."""
    import requests
    from xml.etree import ElementTree
    result = {}
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            state = tree.findtext('state', default='stopped')
            file_name = None
            file_xpath_list = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
            for xp in file_xpath_list:
                elem = tree.find(xp)
                if elem is not None and elem.text:
                    file_name = os.path.basename(elem.text)
                    break
            result['state'] = state
            result['file_name'] = file_name
        else:
            result['state'] = 'unknown'
            result['file_name'] = None
    except Exception as e:
        result['state'] = 'unknown'
        result['file_name'] = None
        logger_qw35sft2_733be7.warning(f'VLC status fetch error: {e}')
    try:
        cmd_result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")
        config_path = cmd_result.get('output', '').strip()
    except Exception:
        config_path = '/home/user/.config/vlc/vlcrc'
    content = env.controller.get_file(config_path)
    if content:
        config_text = content.decode('utf-8', errors='replace')
        record_path = ''
        for line in config_text.split('\n'):
            stripped = line.strip()
            if stripped.startswith('#') or not stripped:
                continue
            if 'input-record-path=' in stripped:
                record_path = stripped.split('=', 1)[-1].strip()
        result['record_path'] = record_path
    else:
        result['record_path'] = None
    return result

def get_mp3_file_size__55ceccfa8167b4a354049a9775b28527_qw35sft2_07f0402e(env, config: dict):
    """Get the byte size of Baby Justin Bieber.mp3 on the desktop."""
    result = env.controller.run_bash_script('stat -c%s "/home/user/Desktop/Baby Justin Bieber.mp3" 2>/dev/null || echo 0', timeout=10)
    if result is None:
        return {'size': 0}
    output = result.get('output', '0').strip()
    try:
        size = int(output)
    except (ValueError, TypeError):
        size = 0
    return {'size': size}

def get_vlc_maxvol_and_bgcone_expands__d25099f13e7da2513cf33c07be0e7955_qw35sft2_c095702e(env, config: dict):
    """Read qt-max-volume and qt-bgcone-expands from VLC's vlcrc config file."""
    result = {'qt_max_volume': '125', 'qt_bgcone_expands': '1'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_6483e4.warning('vlcrc not found or empty')
        return result
    for line in vlcrc_bytes.decode('utf-8', errors='ignore').split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
        elif 'qt-bgcone-expands=' in line:
            result['qt_bgcone_expands'] = line.split('=', 1)[-1].strip()
    logger_qw35sft2_6483e4.info('VLC maxvol+bgcone-expands state: %s', result)
    return result

def get_vlc_cone_and_minview__24dc6e88ff00c69857d313c4738a8f31_qw35sft2_e221381a(env, config: dict):
    """Read VLC config and return qt-bgcone and qt-minimal-view settings."""
    try:
        result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")
        config_path = result.get('output', '').strip()
    except Exception:
        config_path = '/home/user/.config/vlc/vlcrc'
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'vlcrc not found'}
    config_text = content.decode('utf-8', errors='replace')
    qt_bgcone = '1'
    qt_minimal_view = '0'
    for line in config_text.split('\n'):
        if 'qt-bgcone=' in line:
            qt_bgcone = line.split('=')[-1].strip()
        if 'qt-minimal-view=' in line:
            qt_minimal_view = line.split('=')[-1].strip()
    return {'qt-bgcone': qt_bgcone, 'qt-minimal-view': qt_minimal_view}

def get_vlc_record_path__f613292e2de2a0e0145aaa24fdebbc72_qw35sft2_efd6fb7c(env, config: dict):
    """Get VLC recording path from vlcrc configuration file."""
    try:
        result = env.controller.run_bash_script("grep -E '^record-path=' /home/user/.config/vlc/vlcrc 2>/dev/null | head -1", timeout=10)
        if result is None:
            return {'record_path': None, 'error': 'No result from bash script'}
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        output = output.strip()
        if output and '=' in output:
            path = output.split('=', 1)[1].strip()
            return {'record_path': path}
        return {'record_path': None}
    except Exception as e:
        return {'record_path': None, 'error': str(e)}

def get_vlc_stream_and_volume__a51ef1127f9b7bededbd54fd205d555d_qw35sft2_dfc2e8ed(env, config: dict):
    """Read VLC recent MRL list and qt-max-volume setting from vlcrc."""
    mrl_content = ''
    try:
        mrl_bytes = env.controller.get_file('/home/user/.config/vlc/vlc-qt-interface.conf')
        if mrl_bytes:
            mrl_content = mrl_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_e86c9a.warning('Could not read vlc-qt-interface.conf: %s', e)
    recent_mrl = ''
    for line in mrl_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('list='):
            recent_mrl = stripped[5:].strip()
            break
    vlcrc_content = ''
    try:
        vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
        if vlcrc_bytes:
            vlcrc_content = vlcrc_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_e86c9a.warning('Could not read vlcrc: %s', e)
    qt_max_volume = '125'
    for line in vlcrc_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('qt-max-volume='):
            qt_max_volume = stripped.split('=', 1)[1].strip()
            break
    return {'recent_mrl': recent_mrl, 'qt-max-volume': qt_max_volume}

def get_vlc_snapshot_pictures__cdcbbd90330cd55be87f2b9353b8c43c_qw35sft2_9917fb41(env, config: dict):
    """Check if scene.png exists in Pictures and is a valid PNG."""
    file_bytes = env.controller.get_file('/home/user/Pictures/scene.png')
    if not file_bytes or len(file_bytes) < 100:
        return {'exists': False, 'is_png': False, 'size': 0}
    is_png = file_bytes[:4] == b'\x89PNG'
    return {'exists': True, 'is_png': is_png, 'size': len(file_bytes)}

def get_vlc_fullscreen_and_bgcone__986ee683b9dd7fb42919cda3330ee515_qw35sft2_b9d5fc00(env, config: dict):
    """Get VLC fullscreen state and the qt-bgcone setting from vlcrc."""
    result = {}
    screen_size = env.controller.get_vm_screen_size()
    window_size = env.controller.get_vm_window_size(app_class_name='vlc')
    if screen_size and window_size:
        result['is_fullscreen'] = window_size.get('width') == screen_size.get('width') and window_size.get('height') == screen_size.get('height')
    else:
        result['is_fullscreen'] = False
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    result['qt_bgcone'] = '1'
    if vlcrc_bytes:
        vlcrc_content = vlcrc_bytes.decode('utf-8', errors='ignore')
        for line in vlcrc_content.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'qt-bgcone=' in line:
                result['qt_bgcone'] = line.split('=', 1)[-1].strip()
                break
    logger_qw35sft2_43a649.info('VLC fullscreen+bgcone state: %s', result)
    return result

def get_vlc_loop_prefs__a02aa930219f8e3a3317ea91cceaa2ff_qw35sft2_52cb7a7f(env, config: dict):
    """Read VLC instance-from-file and loop settings from vlcrc."""
    result = env.controller.run_bash_script("grep -E '^#?(one-instance-when-started-from-file|loop)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
    if isinstance(result, dict):
        text = result.get('output', result.get('stdout', ''))
    else:
        text = str(result) if result else ''
    uncommented = {}
    commented = {}
    for line in text.strip().splitlines():
        line = line.strip()
        if not line:
            continue
        is_comment = line.startswith('#')
        raw = line[1:].strip() if is_comment else line
        if '=' in raw:
            key, _, val = raw.partition('=')
            key = key.strip()
            val = val.strip()
            if is_comment:
                commented[key] = val
            else:
                uncommented[key] = val
    merged = {**commented, **uncommented}
    return {'one_instance_from_file': merged.get('one-instance-when-started-from-file', 'not_set'), 'loop': merged.get('loop', 'not_set')}

def get_vlc_status_with_loop__c8546449445cd5e2551a333da05cf053_qw35sft2_dfccd919(env, config: dict):
    """Get VLC playing status and loop mode from the HTTP status API."""
    import requests
    from xml.etree import ElementTree
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
    except Exception as e:
        return {'error': str(e)}
    if response.status_code != 200:
        return {'error': f'HTTP {response.status_code}'}
    try:
        tree = ElementTree.fromstring(response.content)
    except Exception as e:
        return {'error': f'XML parse error: {e}'}
    state = tree.findtext('state', default='stopped')
    loop = tree.findtext('loop', default='false')
    file_name = None
    file_xpath_list = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
    for xp in file_xpath_list:
        elem = tree.find(xp)
        if elem is not None and elem.text:
            file_name = os.path.basename(elem.text)
            break
    return {'state': state, 'loop': loop, 'file_name': file_name}

def get_vlc_max_volume__a9382a8b8a5c0cd01bfb6f50454b9a65_qw35sft2_cbb67ad4(env, config: dict):
    """Read qt-max-volume from VLC's vlcrc config file."""
    result = {'qt_max_volume': '125'}
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    if not vlcrc_bytes:
        logger_qw35sft2_d2296b.warning('vlcrc not found or empty')
        return result
    for line in vlcrc_bytes.decode('utf-8', errors='ignore').split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if 'qt-max-volume=' in line:
            result['qt_max_volume'] = line.split('=', 1)[-1].strip()
            break
    logger_qw35sft2_d2296b.info('VLC qt-max-volume: %s', result)
    return result

def get_vlc_cone_and_recordpath__914a6c12022e434d9eb3c729da97ace4_qw35sft2_67033980(env, config: dict):
    """Read VLC config and return qt-bgcone and input-record-path settings."""
    try:
        result = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")
        config_path = result.get('output', '').strip()
    except Exception:
        config_path = '/home/user/.config/vlc/vlcrc'
    content = env.controller.get_file(config_path)
    if not content:
        return {'error': 'vlcrc not found'}
    config_text = content.decode('utf-8', errors='replace')
    qt_bgcone = '1'
    record_path = ''
    for line in config_text.split('\n'):
        if 'qt-bgcone=' in line:
            qt_bgcone = line.split('=')[-1].strip()
        stripped = line.strip()
        if stripped.startswith('input-record-path='):
            record_path = stripped.split('=', 1)[-1].strip()
    return {'qt-bgcone': qt_bgcone, 'input-record-path': record_path}

def get_mp3_existence__77b822a22aadd202cc6c547979661a88_qw35sft2_ad415aff(env, config: dict):
    """Check if Baby Justin Bieber.mp3 exists on the desktop."""
    file_bytes = env.controller.get_file(MP3_PATH_qw35sft2_a7fba1)
    exists = file_bytes is not None and len(file_bytes) > 0
    return {'exists': exists}

def get_vlc_snap_on_desktop__dc35efec6c78cc3245c9e38e3e5c5d4a_qw35sft2_82878c87(env, config: dict):
    """Check if a vlcsnap-* PNG file exists on Desktop and VLC snapshot-path is configured."""
    snap_result = env.controller.run_bash_script('ls /home/user/Desktop/vlcsnap-*.png 2>/dev/null | head -1', timeout=10)
    if isinstance(snap_result, dict):
        snap_output = snap_result.get('output', snap_result.get('stdout', ''))
    else:
        snap_output = str(snap_result)
    has_snap = bool(snap_output.strip())
    pref_result = env.controller.run_bash_script('grep "snapshot-path" /home/user/.config/vlc/vlcrc 2>/dev/null | head -1', timeout=10)
    if isinstance(pref_result, dict):
        pref_output = pref_result.get('output', pref_result.get('stdout', ''))
    else:
        pref_output = str(pref_result)
    path_is_desktop = 'Desktop' in pref_output or '/home/user/Desktop' in pref_output
    return {'snap_on_desktop': has_snap, 'snap_path_configured': path_is_desktop}

def get_wallpaper_and_vlc_running__2cdc858f118b65a32b6755071b250267_qw35sft2_22fb9858(env, config: dict):
    """Get wallpaper URI and check if VLC is still running."""
    wall_cmd = 'gsettings get org.gnome.desktop.background picture-uri'
    wall_result = env.controller.run_bash_script(wall_cmd, timeout=15)
    if isinstance(wall_result, dict):
        wall_output = wall_result.get('output', wall_result.get('stdout', ''))
    else:
        wall_output = str(wall_result)
    wallpaper_uri = wall_output.strip().strip('\'"')
    vlc_cmd = 'pgrep -x vlc | wc -l'
    vlc_result = env.controller.run_bash_script(vlc_cmd, timeout=15)
    if isinstance(vlc_result, dict):
        vlc_output = vlc_result.get('output', vlc_result.get('stdout', '0'))
    else:
        vlc_output = str(vlc_result)
    try:
        vlc_running = int(vlc_output.strip()) > 0
    except ValueError:
        vlc_running = False
    return {'wallpaper_uri': wallpaper_uri, 'vlc_running': vlc_running}

def get_vlc_fullscreen_and_recfolder__27e3d8a3f9225f7e9bc7da9d0f1debd6_qw35sft2_6eeae1b7(env, config: dict):
    """Get VLC fullscreen state and the configured recording folder from vlcrc."""
    result = {}
    screen_size = env.controller.get_vm_screen_size()
    window_size = env.controller.get_vm_window_size(app_class_name='vlc')
    if screen_size and window_size:
        result['is_fullscreen'] = window_size.get('width') == screen_size.get('width') and window_size.get('height') == screen_size.get('height')
    else:
        result['is_fullscreen'] = False
    vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
    result['input_record_path'] = None
    if vlcrc_bytes:
        vlcrc_content = vlcrc_bytes.decode('utf-8', errors='ignore')
        for line in vlcrc_content.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'input-record-path=' in line:
                result['input_record_path'] = line.split('=', 1)[-1].strip()
                break
    logger_qw35sft2_b81e3a.info('VLC fullscreen+recfolder state: %s', result)
    return result

def get_vlc_stream_and_ontop__5d8c2e81cb9645b71051b528d54492c6_qw35sft2_565e3d78(env, config: dict):
    """Read VLC recent MRL list and video-on-top setting from vlcrc."""
    mrl_content = ''
    try:
        mrl_bytes = env.controller.get_file('/home/user/.config/vlc/vlc-qt-interface.conf')
        if mrl_bytes:
            mrl_content = mrl_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_c1f80a.warning('Could not read vlc-qt-interface.conf: %s', e)
    recent_mrl = ''
    for line in mrl_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('list='):
            recent_mrl = stripped[5:].strip()
            break
    vlcrc_content = ''
    try:
        vlcrc_bytes = env.controller.get_file('/home/user/.config/vlc/vlcrc')
        if vlcrc_bytes:
            vlcrc_content = vlcrc_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger_qw35sft2_c1f80a.warning('Could not read vlcrc: %s', e)
    video_on_top = '0'
    for line in vlcrc_content.split('\n'):
        stripped = line.strip()
        if stripped.startswith('video-on-top='):
            video_on_top = stripped.split('=', 1)[1].strip()
            break
    return {'recent_mrl': recent_mrl, 'video-on-top': video_on_top}

def get_vlc_prefs__a7d5b4b88d0304ec4fb07fa10130fbc0_qw35sft2_fbd3b038(env, config: dict):
    """Get VLC recording path and video-title-show preference from vlcrc."""
    try:
        result = env.controller.run_bash_script("grep -E '^(record-path|video-title-show)=' /home/user/.config/vlc/vlcrc 2>/dev/null", timeout=10)
        if result is None:
            return {'record_path': None, 'video_title_show': None, 'error': 'No result from bash script'}
        output = result.get('output', '') if isinstance(result, dict) else str(result)
        record_path = None
        video_title_show = None
        for line in output.strip().splitlines():
            line = line.strip()
            if line.startswith('record-path='):
                record_path = line.split('=', 1)[1].strip()
            elif line.startswith('video-title-show='):
                video_title_show = line.split('=', 1)[1].strip()
        return {'record_path': record_path, 'video_title_show': video_title_show}
    except Exception as e:
        return {'record_path': None, 'video_title_show': None, 'error': str(e)}
