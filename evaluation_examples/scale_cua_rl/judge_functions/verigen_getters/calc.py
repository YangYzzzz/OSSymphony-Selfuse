"""VeriGen RL judge functions.

Source: getters.py
This module is auto-split from the original merged rl_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile
from datetime import date
from datetime import timedelta, date
import base64
import tempfile, os

logger = logging.getLogger(__name__)
logger_qw35sft2_a8fd39 = logging.getLogger(__name__)
logger_qw35sft2_407382 = logging.getLogger('desktopenv.getters.gimp_custom')
_NS_TEXT_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_7af95d = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
_NS_LOEXT_qw35sft2_7af95d = 'urn:org:documentfoundation:names:experimental:office:xmlns:loext:1.0'
_HIGHLIGHT_TRANSPARENT_qw35sft2_7af95d = {'transparent', '', 'automatic', '#00000000', 'none'}
_NS_TEXT_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:text:1.0'
_NS_STYLE_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:style:1.0'
_NS_FO_qw35sft2_fd94cd = 'urn:oasis:names:tc:opendocument:xmlns:xsl-fo-compatible:1.0'
logger_qw35sft2_a4e44f = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_fff6f2 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_9ad635 = logging.getLogger('desktopenv.getters.os')
logger_qw35sft2_a25b79 = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_2c81bb = logging.getLogger('desktopenv.getters.os_custom')
logger_qw35sft2_edeb6d = logging.getLogger('desktopenv.getters.eml_backup_state')
logger_qw35sft2_7606d8 = logging.getLogger(__name__)
logger_qw35sft2_c04067 = logging.getLogger(__name__)
logger_qw35sft2_f39cbb = logging.getLogger('desktopenv.getters.eml_count__2731b9abd5cfbad9ed4df8aae737addc')
logger_qw35sft2_848a1e = logging.getLogger(__name__)
logger_qw35sft2_0d8aba = logging.getLogger('desktopenv.getters.eml_listing')
logger_qw35sft2_e8ef71 = logging.getLogger(__name__)
logger_qw35sft2_fd9ea1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_fd9ea1 = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_05799f = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c7a48d = logging.getLogger(__name__)
logger_qw35sft2_6b91bf = logging.getLogger(__name__)
logger_qw35sft2_a330eb = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_78b50b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_d3002b = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_d3002b = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_5215b8 = logging.getLogger(__name__)
logger_qw35sft2_a85c07 = logging.getLogger('desktopenv.getters.vlc_traj_verify_1')
logger_qw35sft2_2ae123 = logging.getLogger(__name__)
logger_qw35sft2_bf2b87 = logging.getLogger(__name__)
logger_qw35sft2_a3ebf5 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_0975e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_bfc3c0 = logging.getLogger('desktopenv.getters.vlc_traj_verify_4')
logger_qw35sft2_0f05b1 = logging.getLogger(__name__)
logger_qw35sft2_4fe03f = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp3'
MP4_PATH_qw35sft2_4fe03f = '/home/user/Desktop/Baby Justin Bieber.mp4'
logger_qw35sft2_ebcb4c = logging.getLogger(__name__)
logger_qw35sft2_98c532 = logging.getLogger(__name__)
TARGET_URL_qw35sft2_98c532 = 'https://devstreaming-cdn.apple.com/videos/streaming/examples/img_bipbop_adv_example_fmp4/master.m3u8'
logger_qw35sft2_8bd001 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_733be7 = logging.getLogger(__name__)
logger_qw35sft2_8fe0d5 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_8fe0d5 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_6483e4 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_a9c05e = logging.getLogger(__name__)
logger_qw35sft2_e86c9a = logging.getLogger(__name__)
logger_qw35sft2_43a649 = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_fc28df = logging.getLogger(__name__)
logger_qw35sft2_d2296b = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_31805d = logging.getLogger(__name__)
logger_qw35sft2_a7fba1 = logging.getLogger('desktopenv.getters.vlc_custom')
MP3_PATH_qw35sft2_a7fba1 = '/home/user/Desktop/Baby Justin Bieber.mp3'
logger_qw35sft2_b81e3a = logging.getLogger('desktopenv.getters.vlc_traj')
logger_qw35sft2_c1f80a = logging.getLogger(__name__)
logger_qw35sft2_fa1173 = logging.getLogger(__name__)
logger_qw35sft2_c93c96 = logging.getLogger(__name__)
logger_qw35sft2_4c210b = logging.getLogger(__name__)
logger_qw35sft2_d1773d = logging.getLogger(__name__)
logger_qw35sft2_2620a4 = logging.getLogger(__name__)
logger_qw35sft2_43fd17 = logging.getLogger(__name__)
logger_qw35sft2_5d3c04 = logging.getLogger(__name__)
logger_qw35sft2_2ed96e = logging.getLogger(__name__)
logger_qw35sft2_c6cea1 = logging.getLogger(__name__)
logger_qw35sft2_b0d92d = logging.getLogger(__name__)
logger_qw35sft2_65fcf1 = logging.getLogger(__name__)
logger_qw35sft2_52d4f8 = logging.getLogger(__name__)

__all__ = ['get_xlsx_column_a_names__a68724700bb60550a33fa328934ea4ee', 'get_xlsx_row_scores__1700c01d5ad41ab8401de0bafb801990', 'get_xlsx_sheet_names__64efe2d9a5895dad59edc0eb16307294', 'get_xlsx_cell_a21__dc007ee2aea5c445c94303ee0d01fcf4', 'get_xlsx_sheet_names__1a11ade1a8ab5d9c5b9f1e6b66e05fc9', 'get_xlsx_column_c_data__bef74e7adbdaaf2c559cbe5b3bf80940', 'get_xlsx_period_rate_col__02a40d6cc593a8648c91f6def88b5342', 'get_xlsx_initials_column__06e4265c097251cdade138725012bf19', 'get_xlsx_row_values__b46e00d9bcf11154ef195dbbf421e737', 'get_xlsx_cell_value__9956b78e7ee895cd5df580471d77a6ec', 'get_xlsx_cell_value__3f903c88c3cac4cb7f6f63b58e6790e4', 'get_xlsx_cell_h2_value__3ad3a7a9e15d5bbcab97b429c7a2d70d', 'get_xlsx_multi_cell_values__f0c5ce9b9ab74cdf70d3ea75ec500ce3', 'get_xlsx_col_d_text__598fab012ada3ac788280b75e5c01391', 'get_xlsx_multi_cell_values__7424fba2d2a12056c33da3db49551203', 'get_xlsx_multi_cells__1cd51ca9c28aca9840aa590661974218', 'get_xlsx_email_cells__93729a4c1a924bf8a2b603d4f55d7635', 'get_xlsx_range_cells__980775095293f828bec4466da30db27b', 'get_xlsx_sum_cell__b86ee128372bf8371e6a9412924e6faf', 'get_xlsx_cell_value__3aaf5eb76a2414410378f220bee2678f', 'get_xlsx_column_values__312ae9a7a4fdef961589d00f7ec76217', 'get_xlsx_cell_value__c13e5303d821a816dfe6dbe089834e52', 'get_xlsx_cell_value__fe91ae8d0e7c34f2eef1f4db512bed2a', 'get_xlsx_row_by_label__0cd3deb0e30597bee67bd1ee61842fcf', 'get_xlsx_year_cells__dcb80f39266c26be10a70749c6e4e4f7', 'get_xlsx_column_values__0c918bbe7b80a9e03be9fde4824c7824', 'get_xlsx_cell_value__f782061fb40a07e3f7bb2ac7aceef53a', 'get_xlsx_cell_value__d0a50343ffaff2e09611bc0e4f880cee', 'get_xlsx_cell_value__d4eda5c6ae488cf666c27c91b5fbd879', 'get_xlsx_sheet2_merge_pattern__3f03b2d5555a7e222477184511a0b152', 'get_xlsx_cell_value__f3f19b96be4e433759cb0f99e36be750', 'get_xlsx_cells_v0__c8482c50afd0db1eeb95cd8ec431bb1e', 'get_xlsx_column_values__6d4300c0b784865a64e946f955774003', 'get_xlsx_column_order__5f28938c5d9137ac59ddda20f23ca5f2', 'get_xlsx_row_data__c6556bc153411ec2613fae626891d4a6', 'get_xlsx_sorted_titles__4841f866c3a0d68d763e22faff3652cb', 'get_xlsx_cell_values__456143f49bf34353a5e90a3c287b15a6', 'get_xlsx_excellent_counts__4d92f18ea2c033e230b5011ee11fad9e', 'get_xlsx_header_values__79ab7a7fafe67a8a56664d433902c4eb', 'get_xlsx_cell_value__774d1bad2b697745d27771510730ed1b', 'get_xlsx_year_column__5d33fe5f963081b2f481ac5f73c3ccbc', 'get_xlsx_column_values__4ca76e68b8033d8d7fa4f889ee4aceef', 'get_xlsx_sheet_name__e743a4ede9bcdbd814ac13f073acba4a', 'get_xlsx_header_cells__25536421a644aed93ef4482447c70e86', 'get_xlsx_header_row__51f83653b8f37d64288aaccfa8443c8f', 'get_xlsx_cell_value__269bac677b4d800ce51686bd69d7a5b3', 'get_xlsx_column_values__3201def6075c63c69d0100c10729d61d', 'get_xlsx_cells_c8_c11__6709fbc098f9a1f484c1ad8d161e48b0', 'get_xlsx_column_values__a9d0754cd398d25ba1c2a300dbecfcf0', 'get_xlsx_cell_value__da67ff83a4f2f562e6cd5555f90c78ad', 'get_xlsx_cell_value__7aae35559065c2f6c82a9e66b7d55095', 'get_xlsx_cells_zone1_row_totals__b15b7a7665c71e3a9519658d5a13d57f', 'get_xlsx_cell_value__13d011eac188a961b7fc2a44b3a2069f', 'get_xlsx_row_scores__dfb627af00a19217031ee997d3cc240d', 'get_xlsx_two_cells__5caee571623589595fc08cd4d0ebf084', 'get_xlsx_cell_value__02607caca98b58ee759729206fe827f7', 'get_xlsx_cell_value__11373fb010a9f0df0f8a282c8597081e', 'get_xlsx_multi_cell__1393f3a73542b23a8e048e3ef0292819', 'get_xlsx_two_cells__8918eac5e4b67ba092cf5e7258979923', 'get_timetable_cell__9319f3dc4336ee9ea3553360571f0046', 'get_xlsx_sheet2_units_by_product__8ec7487c1752ceb36e17eb91a13203df', 'get_xlsx_cell_value__04a818774a969e111c0b34c8156dff9c', 'get_xlsx_cell_value__bbbcaf548e27f0a1a85c1f51c758e3f5', 'get_xlsx_cell_value__a79441059a79951e2976eaf683658a71', 'get_xlsx_chart_info__9affe8a59575f48a48ea3570040f007c', 'get_xlsx_cell_value__85f75e16b65eecbecc9a6c1bce611233', 'get_xlsx_cell_value__2b6911a42a0fb61cba672fa8735595e9', 'get_xlsx_sheet_names__424814124c3140f33ad68e6e88a912a8', 'get_xlsx_cells_b3_b6__f927596325f5c4e25b12a463bfafbace', 'get_xlsx_cell_value__b28ddd95dc90fbf8e65579083da3a568', 'get_xlsx_cell_value__cb33c763b3ef0ae13364585aa1e6bfa3', 'get_xlsx_profit_column__11a8fd5475e2aba6ae95fda21befab56', 'get_xlsx_sheet2_layout__8c3665733cda1e13a57dc5b562c81969', 'get_xlsx_cell_value__5729380214ac80099631e898f8c00bcc', 'get_xlsx_header_value__e0d3491cdcd2d9a19af59b5d2c376712', 'get_xlsx_row_scores__ba349a57c83f48325e3b8e8c17891359', 'get_xlsx_new_row__246ef0821ed926b9e34f1369cfdfe795', 'get_xlsx_earliest_paper__53a0fb3c44861b74092bea203bb878cc', 'get_xlsx_cell_value__db45e57fbcc1253e5b0f6029b5341c3e', 'get_xlsx_cell_value__4a0f3185fd3c30bf019ccadcf5573a11', 'get_xlsx_cell_value__488a286f36e304dc5e1f93ce39ded8cc', 'get_xlsx_column_hidden__a4f4a609a17cec62166f8b2d532de58f', 'get_xlsx_cells_zone3_row_totals__c68f755a470f57d0018a142ca4ed4f38', 'get_xlsx_zoom__9fd4567c6d7b24fec1b69292feb79e7d', 'get_xlsx_cell_bgcolor__2562d5d5c00c43d3ebb3cc2320483fc3', 'get_xlsx_total_row__6a822179da68f351e1e0a8c5b6636775', 'get_xlsx_summary_rep_totals__8a3eb44c98a42fd6dcd94122bb8b7a9a', 'get_xlsx_column_values__40371a0bf41fabc605f9425ce4121622', 'get_xlsx_column_values__aad13293a214d4a4614be1039d459c78', 'get_xlsx_zoom_and_freeze__048259c26ab013ccdf7c1dac5c5ed24c', 'get_xls_cell_value__2674782c2441965e8475070c5235e39b', 'get_xlsx_cell_value__936e4dec15bc5593ae1b9d4f9d1cf524', 'get_xlsx_sheet2_month_avg__d92436572f728d9c9f2cb1a9c19eb76d', 'get_xlsx_multi_cell__6d0d805c4a52313050aa161e6524aaf8', 'get_xlsx_sheet2_column_data__8a87fce5953091d60c1165596b42e521', 'get_xlsx_sheet_data__7192ac900176bf2dc148847f36558198', 'get_xlsx_row_by_label__19698d046d0da22c68c3bfcbfbd57b88', 'get_xlsx_cell_value__0ba689dbed2285037ec7f44756763cc8', 'get_xlsx_profit_margin_col__60b13be5a7ccd193b76ecf097a198265', 'get_xlsx_cell_bgcolor__fabb4ab7d265226ac33712a7dcf2e535', 'get_xlsx_multi_column_cells__9c6205dd9e0fb3145a3021bfaf014d23', 'get_xlsx_column_header_and_values__23bfc830003e8f489b53a39808f56b59', 'get_xlsx_cell_value__3440d9b31b1654ffe8da4da80be4e5fc', 'get_xlsx_cell_value__52b5a1c85fa3a63d940f3aeb966674a8', 'get_xlsx_cell_value__3d7517c3edf35996bcb8466dea369ae2', 'get_xlsx_cell_value__8ed2a3a912f306294d5ece5e575ff288', 'get_xlsx_cumulative_col__6bc8ae788974152807b9723497a26da7', 'get_xlsx_column_values__34bc68338cfc2efa5ed7f80b3fe4afc3', 'get_xlsx_multi_cells__677226b05224be2ecbd2de2e693f9c8e', 'get_xlsx_cell_value__5eff620ab5e8861fd0f5d3b257b5ca40', 'get_xlsx_row_values__e0472552b507a6465c494b007b9b0305', 'get_xlsx_column_values__2b83d943849f7c327e47bb940fe75967', 'get_xlsx_sum_row__99d61f333e953d709c5fe5f221bfb63e', 'get_timetable_cell__929c2cac5efeaad2cf5e1556b5f1f38f', 'get_xlsx_names_and_city__95be0722e22941aff21eff3d680c0c01', 'get_xlsx_cell_value__da4a3fddcbcd3738d01b04b1ba353fc4', 'get_xlsx_column_values__56ea06753821989747521678ff3594b3', 'get_xlsx_cells__bf5b353ae5739d8c5c981c28dcf67581', 'get_xlsx_cell_value__3a7517792dc1ec4e3be4f857bdb0950f', 'get_xlsx_sheet_data__0d0f959234d0c494b6d7c17cb5158836', 'get_xlsx_cells_zone2_col_totals__dfa00147f689f96c837b3a22754c016b', 'get_xlsx_cell_value__3dfc90e20421740ec2b3449751e1ede4', 'get_xlsx_sorted_sales__38a9aeefb36e58a6a0dcece70b9af3ed', 'get_xlsx_clevel_count__802aa06dfc0f94fbf79660d5d540798f', 'get_xlsx_cell_value__27fb6350ba8922eab313e55ef2751ee9', 'get_xlsx_column_values__0762adbae4ad235c185f47096cf64c91', 'get_xlsx_column_values__40424b53798c2f42a45c9b5d1f78c9a0', 'get_xlsx_sheet2_sales_by_rep__5c7590b2fce2859cb43c07916cdde25c', 'get_timetable_cell__ae3105aee7e1d6920d1ca91156ffa145', 'get_xlsx_cell_value__9e8bcb68d0493b0f9f821969711fa470', 'get_xlsx_cell_value__8afea61a57c3ecda7ab6098e540526c4', 'get_xlsx_reversed_names__d8dd1c41700e60d858e6fd9c2f1f5451', 'get_xlsx_column_d_values__d026dfbc396400238630feee90a0bfc7', 'get_xlsx_cell_value__02d2303f1946e73506bf1a1871c3326e', 'get_xlsx_sheet_names__7b9538bff9f96ac19c0aebd38f8a5f3f', 'get_xlsx_sheet2_maxmin__56726039aefd1f0c57475732eb3ab1f0', 'get_pptx_table_cell__1546f90c5e4d91a5767b01e5f0a56119', 'get_xlsx_freeze_pane__7c0a43d93f3ca3bc0e18377c9b4b22b9', 'get_xlsx_column_values__542fc695553e3e27095b27327032df35', 'get_xlsx_sheet_data__c2e2ad1d381d1d27297ba9392b0d9cea', 'get_xlsx_column_values__427c1b1ebceb5c2ead27edb218ee7bf9', 'get_xlsx_cell_value__12cf17c32a6915e1266996d25d563e2a', 'get_xlsx_row_values__8c5594571c596e8d15ba168f23066d41', 'get_xlsx_top_performers__85190e1339a2cb1782e74f01270bc172', 'get_xlsx_row_data__00c370b29bbdeefb624743c61a277c9e', 'get_xlsx_row_values__767f7417cb0fc080051f9d7103e674cb', 'get_ods_to_xlsx_check__d8b4b0a7ce586e5e25f592627cc186b5', 'get_xlsx_cell_value__4eed649bcee7f8431b3567000e60f20f', 'get_xlsx_col_e_values__c2b3326804f7e34493f7827960c36c6d', 'get_xlsx_cs_data__6d4ad198dfe8ff0251a084115fc1444c', 'get_xlsx_two_cells__009b0a2429b483dc78b74a0168744c1a', 'get_xlsx_column_values__0f45568bd733d9527d872ad7c053c2c8', 'get_xlsx_cell_value__67a1f86b27a24adb60e8a183c35b0852', 'get_xlsx_cell_value__4abf9f9f6bbbc1d8a7c4b9c7e898d79c', 'get_xlsx_sort_check__08b8e71e292866610cba11baa24dc9d9', 'get_xlsx_cell_value__990c97f8e62011f017525e7c1f376ba9', 'get_xlsx_cell_value__ad5647b3d1b47a9e1996a6e69ca9562d', 'get_xlsx_cell_value__28ae1ff6307904e177d5ad2de88ea142', 'get_xlsx_cell_value__69f94a2f4d67df544b1a5f49c4c42c6d', 'get_xlsx_sheet_names__9889ac58b2fb5d2c955b2821b86f86b7', 'get_xlsx_cell_number_format__092734e9cc221c6f71e099e209498fb4', 'get_xlsx_cell_bgcolor__0e55cf8187f0e07889b1109cdd0b266f', 'get_xlsx_column_values__6d702aff56732f3b1cb6a80e7f740e70', 'get_xlsx_cell_value__2650b8c230f7272f8553cf2f429cf59a', 'get_xlsx_cell_value__9ff5dc914928ae745b032698ddd720c0', 'get_xlsx_cell_b3__0698f7f5da73d25e009ccfe66a594e11', 'get_xlsx_cell_value__2fc57009f3993d012dcfddb9200ab322', 'get_sheet2_merge_text__221a8bb1dbb77d2fdb6346ea5a7576e4', 'get_xlsx_cell_values__3d93b5527badb1dca77f44852a74fbe3', 'get_xlsx_cell_value__949a54a656a59ff2d23607cb949dacfc', 'get_xlsx_multi_cell__8a05426d089eda14388178863fc78e0b', 'get_xlsx_cell_value__0ff3d8d538e7cd587b08458b9c790695', 'get_xlsx_sorted_names__56c311ddff224f797b6614d5151a3c35', 'get_xlsx_cell_value__cb5b990d93f87d0f04a6150ee1d23a05', 'get_xlsx_range_values__160ec07bd1e641be9b31cf92f3848b19', 'get_xlsx_cell_value__d5a8ae815c3da82f8081ee65f4e2cea1', 'get_xlsx_all_data__d9fdc5409ef8a872907b1ddc3db9cd54', 'get_xlsx_number_format__6f34d4ab7744f84c76fe24f84ac55b09', 'get_xlsx_half_rate_col__fd58351f5fafe47d305d3595c5202dcd', 'get_xlsx_sheet2_avg__843226a2563502226c076bf6effbabf5', 'get_xlsx_net_income_column__bc1bb048c0eeb30a1fa5f9604d4f0a11', 'get_xlsx_sort_order__db00ed878507fde13aefe6b71212bcbd', 'get_xlsx_row_by_label__4551fb3f098ec9fe882ac83568cd431f', 'get_xlsx_cell_value__54e9a68f9afb8559eff309088793c132', 'get_xlsx_column_values__70ea2799be65bb426ef2e5f3f76ece43', 'get_xlsx_cell_value__fec4261d3cbd8743e0d3445b3fa0111f', 'get_xlsx_qty_sort_check__d422460342c18f4ab4e35ab11ffce7d8', 'get_xlsx_col_c_values__85a76f9ca092376f5f612fc019e6af39', 'get_xlsx_cell_value__d3f5a4ed5eae2c51d540e2a10864dc43', 'get_xlsx_cell_value__ab8a45e026223c32b7b22c741f1bc3ca', 'get_xlsx_fv_column__5c1e2dfe8bffbdb20d88aa7499677c46', 'get_xlsx_cell_value__e7bff0bae84b8431a03877cfc4fe4751', 'get_xlsx_column_cells__4065e6a4b9dbccc0390b8c5a38e2d1b7', 'get_sheet_names__a164952ae6b41142a59183faf6bedadf_qw35sft2_0b7b29d8', 'get_xlsx_merged_cell_state__48a88b5bb3b362cc9c78cc671cedfc7c_qw35sft2_80edd89c', 'get_xlsx_cell_value__0762adbae4ad235c185f47096cf64c91_qw35sft2_84cd3d68', 'get_sheet2_state__68824c07a6606ac16add37bf4765401b_qw35sft2_ca181f23', 'get_seqno_and_sheet_name__83227706efd58da00bb8213cee393c7a_qw35sft2_a23bfda1', 'get_xlsx_cell_value__cfba55e4566ed0373f9b47631b661201_qw35sft2_b5a6828b', 'get_calc_len_formula__29aa3f16a2f40add619bb0073f57ed2e_qw35sft2_2fc070d1', 'get_xlsx_level_secondary__f558bfbb444b8f5736dbe2588385d5fe_qw35sft2_c0e794c7', 'get_xlsx_cell_value__d4b207925bd2d0c3ababe1e319fa8b1f_qw35sft2_70ffb6dd', 'get_calc_cell_e3__73e90374b4fd7034756570a58d380c35_qw35sft2_d5812f2b', 'get_xlsx_cell_value__1ed2b6f96cb9e3177c42c0765213bb8c_qw35sft2_a3c8bc99', 'get_xlsx_multi_cell__3b5b28cd38bc58813bd97e9b2f5d4ad4_qw35sft2_ebc90768', 'get_xlsx_avg_age_cell__c483fe0106f1e0b2430ad19cd97c77c0_qw35sft2_d43a6ec6', 'get_calc_sheet1_net_income__6142d421fefd6b784b8ad81070a58350_qw35sft2_1a6fd6f4', 'get_xlsx_sort_total_row__2cdc28af4ef71e22984f659197334a50_qw35sft2_15e48de1', 'get_sheet2_sorted_revenue__ff6f4e09739d1959467ff60470ffe2bd_qw35sft2_69d06ae7', 'get_xlsx_cell_b7__9ab5f66fb165e5668bd5b38cac0c73c2_qw35sft2_bb64f4bd', 'get_xlsx_first_col_header__b3e437bf5bf6ec56d9f1c4d46601f3ff_qw35sft2_64bd5de5', 'get_xlsx_cell_colors__03a99dd86c6d30983d2467bc1177489c_qw35sft2_df3a9ab2', 'get_xlsx_transposed_cells__f7593e4503a45a853c1ae96cb08aaf92_qw35sft2_3f7147c6', 'get_calc_sheet_info__ecf996871673b767fb6d252021f29e29_qw35sft2_7e7fd42c', 'get_xlsx_zone2_row_totals__6ef64aee01b9e4fef352e390ba82e0ce_qw35sft2_14275934', 'get_sheet2_three_columns__74fb6bc2ea3e8707d0a1b3dd5202c02d_qw35sft2_f6e61ce4', 'get_calc_c1_d1_state__3be2ab82396c7145b9caded82bed3999_qw35sft2_77ca4cb4', 'get_ramp_accel_cells__f874da4b0492e2e023290240d31f1c8d_qw35sft2_32041586', 'get_sheet2_pct_format__6c9c981a9ca6d475ebfb8637fc973b0b_qw35sft2_4e1f3e8b', 'get_xlsx_padded_max__7fad2dd93f80fb906237dbbf8bdc9fb2_qw35sft2_5fd89774', 'get_xlsx_spent_and_date_format__2b619f85a4e6da743c7b9581de6419ad_qw35sft2_f6d401cf', 'get_calc_pivot_and_sort__6bb297a00c63757f69bb3bc219190d5b_qw35sft2_0f9cdc1c', 'get_xlsx_sheet1_cell__13e2a9cff490c6125c903160099b9b7a_qw35sft2_86b16597', 'get_calc_sort_sum__8b987ab7e49940ce5b75cd3329aaf643_qw35sft2_46f169cb', 'get_xlsx_passfail_and_validation__1231c66fffc58f9761299238f8444bbb_qw35sft2_02b61c7f', 'get_month_pivot_sheet2__360bd21f4036eac9068c6b83aba7e2e2_qw35sft2_faf7f287', 'get_freeze_and_sheet_name__c058d8a3342bf3e7968eb49c80bec94d_qw35sft2_20f6f927', 'get_sheet1_cell_e1__e7e54b6e523b2a67e0ff77298aeb57dc_qw35sft2_41105dd3', 'get_xlsx_cell_value__57a509452fa7aaef388db657e08eeb87_qw35sft2_6006cbc7', 'get_xlsx_cell_value__68e5aa381affaca2d81c1d380acacc93_qw35sft2_b1c6ebcb', 'get_xlsx_cell_value__009b0a2429b483dc78b74a0168744c1a_qw35sft2_3bc353e3', 'get_xlsx_multi_cells__142de785ea8e1e7c66581a5013a34ff3_qw35sft2_c7bebb74', 'get_sheet_names__051a1a86398cf9c50b041fec4af58dd3_qw35sft2_c30ab2c7', 'get_xlsx_cell_value__b0dcd11a1bdcc1e13016c1ab56479e5f_qw35sft2_43b2aa29', 'get_xlsx_cell_c7__76b216d4ab5877594f0cff145afca3ba_qw35sft2_fc72e06d', 'get_calc_summary_sheet__878e408cd8149dcd33226991d9a64c87_qw35sft2_fa464251', 'get_sheet1_net_income_sheet2_revenue__5d938927d392c1d220c2599bb703f167_qw35sft2_9ed717b7', 'get_xlsx_level_primary_range__74bebda311b0ec1306e448288c71ed13_qw35sft2_94755150', 'get_calc_cell_b12__dc5d11c569bf87971ba13d9e911e3365_qw35sft2_1bf8d790', 'get_calc_multi_cells__2412d7535e20b0348d2c25accb946746_qw35sft2_de3f29eb', 'get_xlsx_cell_colors__1a5d60ee39c11d25031d40117b38947f_qw35sft2_ceaed29a', 'get_sheet_names__0c20409273ae865247912386f480b246_qw35sft2_97f98e80', 'get_xlsx_transposed_sorted__f581e5220b9d6fb8293dc4e0f669a755_qw35sft2_eb9648e1', 'get_xlsx_zone1_row_totals__3aa70d99f23178c4e0eb339ad1e35c2e_qw35sft2_5d914765', 'get_sheet2_with_grand_total__0378b78b340823f75ef2df5ebaa9b121_qw35sft2_2dff0f04', 'get_feb_max_cell__663a803da15ba1a41afdeb0114613f5e_qw35sft2_41dfa63d', 'get_calc_bold_header__771ce45281d75f3370a131f92752ad7a_qw35sft2_d69ed0dd', 'get_calc_d1_value__62d484d2d5eca5e0b713e4dbdcc095e8_qw35sft2_08e8e0f0', 'get_sheet2_sorted__2fc0837f7c43e49c220e69070cca644b_qw35sft2_b844fb23', 'get_xlsx_padded_count__50235d7deccc7f57d815dc66d0e983bc_qw35sft2_d9297992', 'get_xlsx_header_bold__636602487d408f3cbef02e77b45f21ff_qw35sft2_c0c36e3c', 'get_xlsx_sort_and_sum__7b14f7b3fd242abdf7c663f7ccb97562_qw35sft2_d057b900', 'get_calc_sort_sumif__bda50c4eee48b4e3a15799921e84da62_qw35sft2_2e4e7995', 'get_xlsx_spent_format_total_row__40a8f46c1e43b9d32b31b2230ffda8db_qw35sft2_fb219596', 'get_xlsx_sales_sum_cell__c61f5322a84515967a3c4cc16d8ae654_qw35sft2_13b519fd', 'get_sheet_name_and_pdf__57d2609933a3290edf4943a154f0921f_qw35sft2_ce5aa5ea', 'get_calc_pivot_and_total__2c779926d5355737435372c9e97f0d18_qw35sft2_43ca166f', 'get_xlsx_passfail_and_count__9a44e82214177e638fd33343975b4139_qw35sft2_82e0b7b8', 'get_pivot_and_sorted__3da7df8645c1252be37e4e927311c6ef_qw35sft2_ee96d6d5', 'get_xlsx_sheet1_cell__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_dfdd8fd9', 'get_xlsx_combined__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_52d26355', 'get_sheet_names__0ed76c34ce42b49c8d0f69a316afbc19_qw35sft2_1c79a8a7', 'get_xlsx_sheet_cell__fd1c3b4eab4669bd9848fe12e4f0d1d4_qw35sft2_48f35820', 'get_xlsx_cells_range__1d24482d98b418ee12818e4d70230d5c_qw35sft2_74c8c004', 'get_xlsx_student_c3_c4__87a420fb18a4347285a9615e2d7a9d87_qw35sft2_70b598e6', 'get_xlsx_sheet_names__08c531e403541f6bab4f59a478d0e6c2_qw35sft2_01a31a75', 'get_sheet_and_cell__ebd5692a1e5abf6c2058f3bfb502b3de_qw35sft2_eb586ff4', 'get_calc_sheet2_3col__581360a1329c03f025e39497ca2b0766_qw35sft2_c19c273c', 'get_calc_multi_cells__ddba4585b2adadf3e6cdf7efce56b822_qw35sft2_bdd7f58f', 'get_xlsx_cell_colors__cc3c6531bd3eb1e9b2189f82dea73399_qw35sft2_64440f9c', 'get_xlsx_cell_value__9d8210ef208924a20f5946d708d3b9b8_qw35sft2_fb0b689c', 'get_pptx_table_cells__d704f09ae59c96034fafb3c1ef256369_qw35sft2_62936050', 'get_docx_table_last_cell__818c616cbd97d0d5bf2ff81918416501_qw35sft2_fcef84bd', 'get_docx_table_with_cell__2aa481b876a15405aedb5c91a8fc78e9_qw35sft2_0860cc06', 'get_xlsx_cell_val__2feaa23241a59b8898b31056b9df1f85_qw35sft2_19216087', 'get_xlsx_freeze_panes__b12f1d02d5521f031b94e1747a1c556a_qw35sft2_1a028055', 'get_xlsx_year_column__c99ee3ea1307143970f7d9489f51bcb4_qw35sft2_99620a62', 'get_xlsx_cell_value__5777a0629a4e670ae915b1fe64e58378_qw35sft2_b06b1277', 'get_xlsx_2017_2018_cities__1272197382489180193ce8b35860bb11_qw35sft2_66402127', 'get_xlsx_cell_val__61d78e261ca58a364d140d5219ceb5c7_qw35sft2_5a873ee6', 'get_xlsx_header_align__7d3ab32fb77b1ea64cfabfdcd56b69aa_qw35sft2_ea7fbc3e', 'get_xlsx_sampled_conf_cities__fb22e54de099b00f4d6a127b12703f11_qw35sft2_e8fe6de6', 'get_xlsx_row_data__5abc6535b144407839fbbe3d8a497678_qw35sft2_983a349c', 'get_docx_and_xlsx__25c4c78ad235f6f7dec2a3ce1f7f2c6f_qw35sft2_6c93d383', 'get_xlsx_cell_val__9d1e0f5a79d29ecfac0c206e85fa98e9_qw35sft2_9aabb6ba', 'get_xlsx_unseen_sheet_headers__ac461f3f5169eddd647d2f7d2aad85a2_qw35sft2_07dfbad5', 'get_xlsx_icml_cities__e73f742a9bd3e72f68e6861886fdd0c7_qw35sft2_1ea5a37e', 'get_xlsx_cell_value__9d415d4db27ab3e1c21d0be924167f0e_qw35sft2_aff163a8']

def get_xlsx_column_a_names__a68724700bb60550a33fa328934ea4ee(env, config: dict):
    """Get restaurant names from column A (A2:A6) of the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        names = []
        for row in range(2, 7):
            val = ws.cell(row=row, column=1).value
            names.append(str(val).strip() if val else None)
        return {'names': names}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_scores__1700c01d5ad41ab8401de0bafb801990(env, config: dict):
    """Get scores from a specific row in the grades spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 3)
        start_col = config.get('start_col', 3)
        end_col = config.get('end_col', 13)
        values = []
        for col in range(start_col, end_col + 1):
            values.append(ws.cell(row=row, column=col).value)
        return {'values': values, 'row': row}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__64efe2d9a5895dad59edc0eb16307294(env, config: dict):
    """Get sheet names from xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        return {'sheet_names': sheet_names}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_a21__dc007ee2aea5c445c94303ee0d01fcf4(env, config: dict):
    """Get cell A21 value from Sheet1."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws.cell(row=21, column=1).value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__1a11ade1a8ab5d9c5b9f1e6b66e05fc9(env, config: dict):
    """Get sheet names from xlsx file in order."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        return {'sheet_names': sheet_names, 'first_sheet': sheet_names[0] if sheet_names else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_c_data__bef74e7adbdaaf2c559cbe5b3bf80940(env, config: dict):
    """Get column C data (header + values) from Sheet1."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws.cell(row=1, column=3).value
        values = []
        for row in range(2, 21):
            val = ws.cell(row=row, column=3).value
            values.append(val)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_period_rate_col__02a40d6cc593a8648c91f6def88b5342(env, config: dict):
    """Get column C header and period rate values from PeriodRate.xlsx."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws['C1'].value
        values = {}
        for row in range(2, 26):
            cell_val = ws.cell(row=row, column=3).value
            if cell_val is not None:
                values[f'C{row}'] = round(float(cell_val), 4)
            else:
                values[f'C{row}'] = None
        return {'header': header, 'values': values, 'count': sum((1 for v in values.values() if v is not None))}
    finally:
        os.unlink(tmp_path)

def get_xlsx_initials_column__06e4265c097251cdade138725012bf19(env, config: dict):
    """Get the initials column (E) and the split data columns (B, C) from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {'header_e': ws.cell(row=1, column=5).value, 'initials': [], 'first_names': [], 'last_names': []}
        for row in range(2, 23):
            result['initials'].append(ws.cell(row=row, column=5).value)
            result['first_names'].append(ws.cell(row=row, column=2).value)
            result['last_names'].append(ws.cell(row=row, column=3).value)
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_values__b46e00d9bcf11154ef195dbbf421e737(env, config: dict):
    """Get values from a specific row in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 12)
        cols = config.get('cols', ['A', 'B', 'C', 'D', 'E', 'F', 'G'])
        result = {}
        for col in cols:
            cell_ref = f'{col}{row}'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__9956b78e7ee895cd5df580471d77a6ec(env, config: dict):
    """Get cell value and check if it's a SUM formula result."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb_data = openpyxl.load_workbook(tmp_path, data_only=True)
        ws_data = wb_data.worksheets[config.get('sheet', 0)]
        cell_ref = config.get('cell', 'C9')
        data_value = ws_data[cell_ref].value
        wb_formula = openpyxl.load_workbook(tmp_path, data_only=False)
        ws_formula = wb_formula.worksheets[config.get('sheet', 0)]
        formula_value = ws_formula[cell_ref].value
        return {'value': data_value, 'formula': str(formula_value) if formula_value else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3f903c88c3cac4cb7f6f63b58e6790e4(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_h2_value__3ad3a7a9e15d5bbcab97b429c7a2d70d(env, config: dict):
    """Read cell H2 from Sheet1 to get the SUMIF result."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws['H2'].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell_values__f0c5ce9b9ab74cdf70d3ea75ec500ce3(env, config: dict):
    """Get multiple cell values from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells = config.get('cells', [])
        values = {}
        for cell in cells:
            values[cell] = ws[cell].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_col_d_text__598fab012ada3ac788280b75e5c01391(env, config: dict):
    """Read column D values (rows 2-30) as text strings from the spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        values = []
        for row in range(2, 31):
            cell = ws.cell(row=row, column=4)
            val = cell.value
            if val is not None:
                values.append(str(val))
            else:
                values.append(None)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell_values__7424fba2d2a12056c33da3db49551203(env, config: dict):
    """Get values of multiple cells from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells_to_read = config.get('cells', [])
        values = {}
        for cell_ref in cells_to_read:
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cells__1cd51ca9c28aca9840aa590661974218(env, config: dict):
    """Get values from multiple specified cells in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cells = config.get('cells', [])
        values = {}
        for cell_ref in cells:
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_email_cells__93729a4c1a924bf8a2b603d4f55d7635(env, config: dict):
    """Get email column values from Professor_Contact.xlsx."""
    import tempfile
    import os
    import openpyxl
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            ws = wb.worksheets[config.get('sheet', 0)]
            cells = config.get('cells', ['F3', 'F4', 'F5'])
            values = {}
            for cell_ref in cells:
                cell_val = ws[cell_ref].value
                values[cell_ref] = str(cell_val).strip() if cell_val is not None else None
            return {'values': values}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading xlsx cells: {e}')
        return {'error': str(e)}

def get_xlsx_range_cells__980775095293f828bec4466da30db27b(env, config: dict):
    """Read a range of cells from an xlsx file including multiple columns and rows."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell_refs = config.get('cell_refs', [])
        cells = {}
        for cell_ref in cell_refs:
            val = ws[cell_ref].value
            cells[cell_ref] = str(val) if val is not None else None
        return {'cells': cells}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sum_cell__b86ee128372bf8371e6a9412924e6faf(env, config: dict):
    """Get the value of a specific cell to check SUM formula result."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'E37')
        value = ws[cell].value
        wb_formula = openpyxl.load_workbook(tmp_path, data_only=False)
        ws_formula = wb_formula.worksheets[config.get('sheet', 0)]
        formula = ws_formula[cell].value
        has_formula = isinstance(formula, str) and formula.startswith('=')
        return {'value': value, 'has_formula': has_formula}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3aaf5eb76a2414410378f220bee2678f(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__312ae9a7a4fdef961589d00f7ec76217(env, config: dict):
    """Get values from a specified column range in an Excel file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        column = config.get('column', 'E')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 10)
        values = []
        for row in range(start_row, end_row + 1):
            cell_ref = f'{column}{row}'
            val = ws[cell_ref].value
            values.append(val)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__c13e5303d821a816dfe6dbe089834e52(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_index = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_index]
        value = ws[cell].value
        return {'value': value, 'cell': cell}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__fe91ae8d0e7c34f2eef1f4db512bed2a(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_by_label__0cd3deb0e30597bee67bd1ee61842fcf(env, config: dict):
    """Scan column A for a label and return values from columns B-G of that row."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        label = config.get('label', 'Max')
        target_row = None
        for row in range(1, ws.max_row + 1):
            cell_val = ws.cell(row=row, column=1).value
            if cell_val and str(cell_val).strip().lower() == label.lower():
                target_row = row
                break
        if target_row is None:
            return {'error': f'Label "{label}" not found in column A', 'label_found': False}
        values = {}
        col_names = ['B', 'C', 'D', 'E', 'F', 'G']
        for (i, col_name) in enumerate(col_names):
            val = ws.cell(row=target_row, column=i + 2).value
            values[col_name] = val
        return {'label_found': True, 'label': str(ws.cell(row=target_row, column=1).value).strip(), 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_year_cells__dcb80f39266c26be10a70749c6e4e4f7(env, config: dict):
    """Get the values of cells B3, B4, B5 to check year entries."""
    file_path = config.get('path', '/home/user/Desktop/best_awards_acl.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'b3': ws['B3'].value, 'b4': ws['B4'].value, 'b5': ws['B5'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__0c918bbe7b80a9e03be9fde4824c7824(env, config: dict):
    """Get values from a specific column range in the xlsx file."""
    import tempfile
    import os
    import openpyxl
    file_path = config.get('path', '/home/user/Resize_Cells_Fit_Page.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 2)
        start_row = config.get('start_row', 12)
        end_row = config.get('end_row', 23)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(str(val) if val is not None else '')
        return {'values': values}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__f782061fb40a07e3f7bb2ac7aceef53a(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__d0a50343ffaff2e09611bc0e4f880cee(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__d4eda5c6ae488cf666c27c91b5fbd879(env, config: dict):
    """Get a specific cell value from a specific sheet in an xlsx file."""
    file_path = config.get('path', '')
    sheet_name = config.get('sheet_name', None)
    cell = config.get('cell', 'A1')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if sheet_name and sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
        else:
            ws = wb.active
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_merge_pattern__3f03b2d5555a7e222477184511a0b152(env, config: dict):
    """Get Sheet2 merge pattern and cell values."""
    file_path = config.get('path', '/home/user/FutureValue.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found'}
        ws = wb['Sheet2']
        merged_cells = set()
        for rng in ws.merged_cells.ranges:
            for cell in rng.cells:
                coord = ws.cell(row=cell[0], column=cell[1]).coordinate
                merged_cells.add(coord)
        result = {}
        for cell_ref in ['A1', 'B1', 'C1', 'D1', 'E1']:
            cell = ws[cell_ref]
            result[cell_ref] = {'value': cell.value, 'merged': cell_ref in merged_cells}
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__f3f19b96be4e433759cb0f99e36be750(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell = config.get('cell', 'A1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_v0__c8482c50afd0db1eeb95cd8ec431bb1e(env, config: dict):
    """Get values of cells A10 and B10 from the Employee Performance Evaluation Summary spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        a10_value = ws['A10'].value
        b10_value = ws['B10'].value
        if isinstance(b10_value, float) and b10_value == int(b10_value):
            b10_value = int(b10_value)
        return {'a10': a10_value, 'b10': b10_value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__6d4300c0b784865a64e946f955774003(env, config: dict):
    """Get values from a specified column range in an Excel file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        column = config.get('column', 'I')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 10)
        values = []
        for row in range(start_row, end_row + 1):
            cell_ref = f'{column}{row}'
            val = ws[cell_ref].value
            values.append(val)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_order__5f28938c5d9137ac59ddda20f23ca5f2(env, config: dict):
    """Get values from a column to verify sort order."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col_letter = config.get('column', 'C')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 8)
        values = []
        for row in range(start_row, end_row + 1):
            cell_val = ws[f'{col_letter}{row}'].value
            values.append(cell_val)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_data__c6556bc153411ec2613fae626891d4a6(env, config: dict):
    """Get cell values from a specific row in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 9)
        result = {}
        for col_letter in config.get('columns', ['A', 'B', 'C', 'D', 'E']):
            cell = ws[f'{col_letter}{row}']
            result[col_letter] = cell.value
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_sorted_titles__4841f866c3a0d68d763e22faff3652cb(env, config: dict):
    """Get titles from column A (rows 2-6) to verify sort order."""
    file_path = config.get('path', '/home/user/Desktop/rsc-ebook-collection-2023.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        titles = []
        for row in range(2, 7):
            val = ws.cell(row=row, column=1).value
            titles.append(val if val else '')
        return {'titles': titles}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_values__456143f49bf34353a5e90a3c287b15a6(env, config: dict):
    """Read values from specified cells."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cells = config.get('cells', [])
        values = {}
        for cell_ref in cells:
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_excellent_counts__4d92f18ea2c033e230b5011ee11fad9e(env, config: dict):
    """Get the header in AA1 and COUNTIF values in AA2:AA8 for Excellent ratings."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws['AA1'].value
        counts = []
        for row in range(2, 9):
            val = ws.cell(row=row, column=27).value
            if isinstance(val, float) and val == int(val):
                val = int(val)
            counts.append(val)
        return {'header': header, 'counts': counts}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_values__79ab7a7fafe67a8a56664d433902c4eb(env, config: dict):
    """Get header values from row 1 of an xlsx file."""
    import openpyxl
    file_path = config.get('path', '/home/user/Desktop/GRF-p5y.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found', 'values': []}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        values = []
        for col in range(1, 5):
            cell_val = ws.cell(row=1, column=col).value
            values.append(str(cell_val).strip() if cell_val is not None else None)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__774d1bad2b697745d27771510730ed1b(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_year_column__5d33fe5f963081b2f481ac5f73c3ccbc(env, config: dict):
    """Get values from column E (rows 1-6) to verify Year column."""
    file_path = config.get('path', '/home/user/Desktop/rsc-ebook-collection-2023.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        values = []
        for row in range(1, 7):
            val = ws.cell(row=row, column=5).value
            values.append(val)
        return {'column_e': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__4ca76e68b8033d8d7fa4f889ee4aceef(env, config: dict):
    """Get all values from a specific column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 3)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            cell_val = ws.cell(row=row, column=col).value
            values.append(cell_val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_name__e743a4ede9bcdbd814ac13f073acba4a(env, config: dict):
    """Get the name of a worksheet from an xlsx file on the VM."""
    file_path = config.get('path', '')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        return {'sheet_name': ws.title}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_cells__25536421a644aed93ef4482447c70e86(env, config: dict):
    """Get values of cells E1 and F1 from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        e1_val = ws['E1'].value
        f1_val = ws['F1'].value
        return {'e1': str(e1_val).strip() if e1_val else None, 'f1': str(f1_val).strip() if f1_val else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_row__51f83653b8f37d64288aaccfa8443c8f(env, config: dict):
    """Get all header values from row 1 of an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        headers = []
        for col in range(1, ws.max_column + 1):
            val = ws.cell(row=1, column=col).value
            headers.append(val)
        first_data_row = []
        for col in range(1, ws.max_column + 1):
            val = ws.cell(row=2, column=col).value
            first_data_row.append(val)
        return {'headers': headers, 'first_data_row': first_data_row, 'num_columns': ws.max_column, 'num_rows': ws.max_row}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__269bac677b4d800ce51686bd69d7a5b3(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__3201def6075c63c69d0100c10729d61d(env, config: dict):
    """Get all values from a specified column in an xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 4)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_c8_c11__6709fbc098f9a1f484c1ad8d161e48b0(env, config: dict):
    """Get the values of cells C8, C9, C10, C11 from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        values = {}
        for row in range(8, 12):
            cell_ref = f'C{row}'
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__a9d0754cd398d25ba1c2a300dbecfcf0(env, config: dict):
    """Read all values from a specified column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 3)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 22)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(str(val).strip() if val is not None else None)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__da67ff83a4f2f562e6cd5555f90c78ad(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__7aae35559065c2f6c82a9e66b7d55095(env, config: dict):
    """Get a cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_zone1_row_totals__b15b7a7665c71e3a9519658d5a13d57f(env, config: dict):
    """Get the Total column (F) values for Zone 1 products (rows 3-5)."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for row in [3, 4, 5]:
            cell_ref = f'F{row}'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__13d011eac188a961b7fc2a44b3a2069f(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_scores__dfb627af00a19217031ee997d3cc240d(env, config: dict):
    """Get scores from a specific row in the grades spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 11)
        start_col = config.get('start_col', 3)
        end_col = config.get('end_col', 13)
        values = []
        for col in range(start_col, end_col + 1):
            values.append(ws.cell(row=row, column=col).value)
        return {'values': values, 'row': row}
    finally:
        os.unlink(tmp_path)

def get_xlsx_two_cells__5caee571623589595fc08cd4d0ebf084(env, config: dict):
    """Get values from two specified cells in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell1 = config.get('cell1', 'A1')
        cell2 = config.get('cell2', 'A2')
        val1 = ws[cell1].value
        val2 = ws[cell2].value
        return {'cell1_value': val1, 'cell2_value': val2}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__02607caca98b58ee759729206fe827f7(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__11373fb010a9f0df0f8a282c8597081e(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell__1393f3a73542b23a8e048e3ef0292819(env, config: dict):
    """Get multiple cell values from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells = config.get('cells', [])
        values = {}
        for cell in cells:
            values[cell] = ws[cell].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_two_cells__8918eac5e4b67ba092cf5e7258979923(env, config: dict):
    """Get values from two specific cells in an xlsx file."""
    file_path = config.get('path', '/home/user/DemographicProfile.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet', 'Sheet1')
        if sheet_name not in wb.sheetnames:
            return {'error': f"Sheet '{sheet_name}' not found"}
        ws = wb[sheet_name]
        cell1 = config.get('cell1', 'A1')
        cell2 = config.get('cell2', 'B1')
        return {'cell1_value': ws[cell1].value, 'cell2_value': ws[cell2].value}
    finally:
        os.unlink(tmp_path)

def get_timetable_cell__9319f3dc4336ee9ea3553360571f0046(env, config: dict):
    """Get cell value and background color from the Course Timetable xlsx."""
    file_path = config.get('path', '/home/user/Desktop/Course Timetable.xlsx')
    cell_addr = config.get('cell', 'F2')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell = ws[cell_addr]
        value = cell.value
        fill = cell.fill
        bg_color = 'none'
        if fill and fill.start_color and fill.start_color.rgb:
            rgb = fill.start_color.rgb
            if rgb and rgb != '00000000':
                bg_color = str(rgb)
        return {'value': value, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_units_by_product__8ec7487c1752ceb36e17eb91a13203df(env, config: dict):
    """Read Sheet2 to get product names and their total units values."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if len(wb.sheetnames) < 2:
            return {'error': 'Sheet2 not found'}
        ws = wb.worksheets[1]
        products = {'Majestic', 'Quad', 'Alpine', 'Carlota', 'Bellen'}
        data = {}
        for row in range(1, ws.max_row + 1):
            for col in range(1, ws.max_column + 1):
                val = ws.cell(row=row, column=col).value
                if val is not None:
                    val_str = str(val).strip()
                    if val_str in products:
                        for c2 in range(1, ws.max_column + 1):
                            if c2 != col:
                                v2 = ws.cell(row=row, column=c2).value
                                if isinstance(v2, (int, float)) and v2 > 0:
                                    data[val_str] = float(v2)
                                    break
        return {'units_by_product': data}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__04a818774a969e111c0b34c8156dff9c(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__bbbcaf548e27f0a1a85c1f51c758e3f5(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__a79441059a79951e2976eaf683658a71(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_chart_info__9affe8a59575f48a48ea3570040f007c(env, config: dict):
    """Get chart information from an xlsx file on the VM, including data series references."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        charts = ws._charts
        chart_count = len(charts)
        chart_types = []
        chart_series_info = []
        for chart in charts:
            chart_type = type(chart).__name__
            chart_types.append(chart_type)
            series_list = []
            for s in chart.series:
                series_entry = {}
                if hasattr(s, 'val') and s.val and hasattr(s.val, 'numRef') and s.val.numRef:
                    series_entry['value_ref'] = s.val.numRef.f
                elif hasattr(s, 'values') and s.values:
                    pass
                if hasattr(s, 'cat') and s.cat and hasattr(s.cat, 'strRef') and s.cat.strRef:
                    series_entry['category_ref'] = s.cat.strRef.f
                elif hasattr(s, 'cat') and s.cat and hasattr(s.cat, 'numRef') and s.cat.numRef:
                    series_entry['category_ref'] = s.cat.numRef.f
                series_list.append(series_entry)
            chart_series_info.append(series_list)
        return {'chart_count': chart_count, 'chart_types': chart_types, 'chart_series_info': chart_series_info}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__85f75e16b65eecbecc9a6c1bce611233(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__2b6911a42a0fb61cba672fa8735595e9(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__424814124c3140f33ad68e6e88a912a8(env, config: dict):
    """Get sheet names from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_b3_b6__f927596325f5c4e25b12a463bfafbace(env, config: dict):
    """Get the values of cells B3, B4, B5, B6 from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        values = {}
        for row in range(3, 7):
            cell_ref = f'B{row}'
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__b28ddd95dc90fbf8e65579083da3a568(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_path = config.get('path', '/home/user/Desktop/researchers.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__cb33c763b3ef0ae13364585aa1e6bfa3(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_profit_column__11a8fd5475e2aba6ae95fda21befab56(env, config: dict):
    """Get header and values from the Profit column (D)."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        header = ws.cell(row=1, column=4).value
        values = []
        for row in range(2, 12):
            values.append(ws.cell(row=row, column=4).value)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_layout__8c3665733cda1e13a57dc5b562c81969(env, config: dict):
    """Get Sheet2 layout: cell values and merge status."""
    file_path = config.get('path', '/home/user/FutureValue.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found'}
        ws = wb['Sheet2']
        merged_cells = set()
        for rng in ws.merged_cells.ranges:
            for cell in rng.cells:
                coord = ws.cell(row=cell[0], column=cell[1]).coordinate
                merged_cells.add(coord)
        result = {}
        for cell_ref in ['A1', 'B1', 'C1', 'D1', 'E1', 'A2', 'B2', 'C2', 'D2', 'E2']:
            cell = ws[cell_ref]
            result[cell_ref] = {'value': cell.value, 'merged': cell_ref in merged_cells}
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__5729380214ac80099631e898f8c00bcc(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_value__e0d3491cdcd2d9a19af59b5d2c376712(env, config: dict):
    """Get the value of cell A1 from the spreadsheet to check header correction."""
    file_path = config.get('path', '/home/user/Desktop/best_awards_acl.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'a1': ws['A1'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_scores__ba349a57c83f48325e3b8e8c17891359(env, config: dict):
    """Get scores from a specific row in the grades spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 9)
        start_col = config.get('start_col', 3)
        end_col = config.get('end_col', 13)
        values = []
        for col in range(start_col, end_col + 1):
            values.append(ws.cell(row=row, column=col).value)
        return {'values': values, 'row': row}
    finally:
        os.unlink(tmp_path)

def get_xlsx_new_row__246ef0821ed926b9e34f1369cfdfe795(env, config: dict):
    """Get values from a specific row in the spreadsheet."""
    import tempfile
    import os
    import openpyxl
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            ws = wb.worksheets[config.get('sheet', 0)]
            row = config.get('row', 6)
            columns = config.get('columns', ['A', 'B', 'C', 'D'])
            values = {}
            for col in columns:
                cell_ref = f'{col}{row}'
                cell_val = ws[cell_ref].value
                if cell_val is not None:
                    values[col] = str(cell_val).strip()
                else:
                    values[col] = None
            return {'row_values': values}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading xlsx row: {e}')
        return {'error': str(e)}

def get_xlsx_earliest_paper__53a0fb3c44861b74092bea203bb878cc(env, config: dict):
    """Get values from cells E1 and E2 to verify earliest paper entry."""
    file_path = config.get('path', '/home/user/Desktop/rsc-ebook-collection-2023.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        e1 = ws['E1'].value
        e2 = ws['E2'].value
        return {'e1': e1, 'e2': e2}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__db45e57fbcc1253e5b0f6029b5341c3e(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__4a0f3185fd3c30bf019ccadcf5573a11(env, config: dict):
    """Get a cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__488a286f36e304dc5e1f93ce39ded8cc(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'C8')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_hidden__a4f4a609a17cec62166f8b2d532de58f(env, config: dict):
    """Check which columns are hidden in the xlsx file."""
    import tempfile
    import os
    import openpyxl
    file_path = config.get('path', '/home/user/Resize_Cells_Fit_Page.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        columns_to_check = config.get('columns', ['E', 'F', 'G', 'H'])
        hidden_status = {}
        for col in columns_to_check:
            hidden_status[col] = ws.column_dimensions[col].hidden
        return {'hidden_status': hidden_status}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_zone3_row_totals__c68f755a470f57d0018a142ca4ed4f38(env, config: dict):
    """Get the Total column (F) values for Zone 3 products (rows 17-19)."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for row in [17, 18, 19]:
            cell_ref = f'F{row}'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_zoom__9fd4567c6d7b24fec1b69292feb79e7d(env, config: dict):
    """Get the zoom scale of the specified sheet in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        zoom_scale = None
        for sv in ws.views.sheetView:
            if sv.zoomScale is not None:
                zoom_scale = sv.zoomScale
                break
        if zoom_scale is None:
            zoom_scale = 100
        return {'zoom': zoom_scale}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_bgcolor__2562d5d5c00c43d3ebb3cc2320483fc3(env, config: dict):
    """Get background colors of specified cells from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cells = config.get('cells', [])
        colors = {}
        for cell_ref in cells:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and fill.fgColor.rgb and (fill.fgColor.rgb != '00000000'):
                colors[cell_ref] = str(fill.fgColor.rgb)
            elif fill and fill.patternType == 'solid' and fill.start_color and fill.start_color.rgb:
                colors[cell_ref] = str(fill.start_color.rgb)
            else:
                colors[cell_ref] = None
        return {'colors': colors, 'count': len(colors)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_total_row__6a822179da68f351e1e0a8c5b6636775(env, config: dict):
    """Get the total row (row 12) values: A12, B12, C12."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        return {'label': ws['A12'].value, 'total_sales': ws['B12'].value, 'total_cogs': ws['C12'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_summary_rep_totals__8a3eb44c98a42fd6dcd94122bb8b7a9a(env, config: dict):
    """Get Sales Rep and Total columns from Summary sheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        target_sheet = None
        for name in wb.sheetnames:
            if name.lower() == 'summary':
                target_sheet = wb[name]
                break
        if target_sheet is None:
            return {'error': 'Summary sheet not found'}
        rep_totals = {}
        for row in range(2, target_sheet.max_row + 1):
            rep_name = target_sheet.cell(row=row, column=1).value
            total_val = target_sheet.cell(row=row, column=2).value
            if rep_name is None:
                break
            rep_totals[str(rep_name).strip()] = total_val
        headers = []
        for col in range(1, target_sheet.max_column + 1):
            val = target_sheet.cell(row=1, column=col).value
            if val is not None:
                headers.append(str(val).strip().lower())
        return {'sheet_found': True, 'headers': headers, 'rep_totals': rep_totals, 'row_count': len(rep_totals)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__40371a0bf41fabc605f9425ce4121622(env, config: dict):
    """Get header and values from a column in an xlsx file."""
    file_path = config.get('path', '')
    header_cell = config.get('header_cell', 'G1')
    column = config.get('column', 7)
    start_row = config.get('start_row', 2)
    end_row = config.get('end_row', 36)
    sheet_index = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_index]
        header = ws[header_cell].value
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=column).value
            values.append(val)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__aad13293a214d4a4614be1039d459c78(env, config: dict):
    """Read a column of values from an xlsx file on the VM."""
    import tempfile
    import os
    import openpyxl
    file_path = config.get('path', '')
    column = config.get('column', 1)
    start_row = config.get('start_row', 1)
    end_row = config.get('end_row', 6)
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'error': 'File not found', 'values': []}
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            ws = wb.worksheets[config.get('sheet', 0)]
            values = []
            for row in range(start_row, end_row + 1):
                cell_val = ws.cell(row=row, column=column).value
                values.append(str(cell_val) if cell_val is not None else '')
            return {'values': values}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e), 'values': []}

def get_xlsx_zoom_and_freeze__048259c26ab013ccdf7c1dac5c5ed24c(env, config: dict):
    """Get both zoom scale and freeze pane setting from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        zoom_scale = None
        for sv in ws.views.sheetView:
            if sv.zoomScale is not None:
                zoom_scale = sv.zoomScale
                break
        if zoom_scale is None:
            zoom_scale = 100
        freeze_panes = ws.freeze_panes
        return {'zoom': zoom_scale, 'freeze_panes': str(freeze_panes) if freeze_panes else None}
    finally:
        os.unlink(tmp_path)

def get_xls_cell_value__2674782c2441965e8475070c5235e39b(env, config: dict):
    """Get a cell value from an .xls file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xls', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        import xlrd
        wb = xlrd.open_workbook(tmp_path)
        ws = wb.sheet_by_index(config.get('sheet', 0))
        row = config.get('row', 0)
        col = config.get('col', 0)
        value = ws.cell_value(row, col)
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__936e4dec15bc5593ae1b9d4f9d1cf524(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_month_avg__d92436572f728d9c9f2cb1a9c19eb76d(env, config: dict):
    """Get Month, Total, and Average columns from Sheet2."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet2 = None
        for name in wb.sheetnames:
            if name.lower() == 'sheet2':
                sheet2 = wb[name]
                break
        if sheet2 is None:
            return {'error': 'Sheet2 not found'}
        headers = []
        for col in range(1, sheet2.max_column + 1):
            val = sheet2.cell(row=1, column=col).value
            if val is not None:
                headers.append(str(val).strip())
        data = {}
        for row in range(2, sheet2.max_row + 1):
            month = sheet2.cell(row=row, column=1).value
            if month is None:
                break
            month = str(month).strip()
            row_data = {}
            for col in range(2, sheet2.max_column + 1):
                header = sheet2.cell(row=1, column=col).value
                if header is not None:
                    row_data[str(header).strip()] = sheet2.cell(row=row, column=col).value
            data[month] = row_data
        return {'headers': headers, 'data': data, 'sheet_found': True}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell__6d0d805c4a52313050aa161e6524aaf8(env, config: dict):
    """Get multiple cell values from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells = config.get('cells', [])
        values = {}
        for cell in cells:
            values[cell] = ws[cell].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_column_data__8a87fce5953091d60c1165596b42e521(env, config: dict):
    """Get column A data (header + values) from Sheet2."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        header = ws.cell(row=1, column=1).value
        values = []
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=1).value
            if val is None:
                break
            values.append(val)
        return {'header': header, 'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_data__7192ac900176bf2dc148847f36558198(env, config: dict):
    """Get all data from a specific sheet in an xlsx file."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', None)
        if sheet_name and sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheets': wb.sheetnames}
        ws = wb[sheet_name] if sheet_name else wb.worksheets[0]
        headers = []
        for col in range(1, ws.max_column + 1):
            headers.append(ws.cell(1, col).value)
        rows = []
        for row in range(2, ws.max_row + 1):
            row_data = {}
            for col in range(1, ws.max_column + 1):
                row_data[headers[col - 1]] = ws.cell(row, col).value
            if any((v is not None for v in row_data.values())):
                rows.append(row_data)
        return {'sheet_name': ws.title, 'headers': headers, 'rows': rows, 'row_count': len(rows), 'all_sheets': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_by_label__19698d046d0da22c68c3bfcbfbd57b88(env, config: dict):
    """Scan column A for a label and return values from columns B-G of that row."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        label = config.get('label', 'Average')
        target_row = None
        for row in range(1, ws.max_row + 1):
            cell_val = ws.cell(row=row, column=1).value
            if cell_val and str(cell_val).strip().lower() == label.lower():
                target_row = row
                break
        if target_row is None:
            return {'error': f'Label "{label}" not found in column A', 'label_found': False}
        values = {}
        col_names = ['B', 'C', 'D', 'E', 'F', 'G']
        for (i, col_name) in enumerate(col_names):
            val = ws.cell(row=target_row, column=i + 2).value
            values[col_name] = val
        return {'label_found': True, 'label': str(ws.cell(row=target_row, column=1).value).strip(), 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__0ba689dbed2285037ec7f44756763cc8(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_profit_margin_col__60b13be5a7ccd193b76ecf097a198265(env, config: dict):
    """Get column D header and values, plus Sales and COGS columns for verification."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        header = ws['D1'].value
        values = []
        for row in range(2, 12):
            val = ws.cell(row=row, column=4).value
            values.append(val)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_bgcolor__fabb4ab7d265226ac33712a7dcf2e535(env, config: dict):
    """Get background colors of specified cells from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cells = config.get('cells', [])
        colors = {}
        for cell_ref in cells:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and fill.fgColor.rgb and (fill.fgColor.rgb != '00000000'):
                colors[cell_ref] = str(fill.fgColor.rgb)
            elif fill and fill.patternType == 'solid' and fill.start_color and fill.start_color.rgb:
                colors[cell_ref] = str(fill.start_color.rgb)
            else:
                colors[cell_ref] = None
        return {'colors': colors, 'count': len(colors)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_column_cells__9c6205dd9e0fb3145a3021bfaf014d23(env, config: dict):
    """Read cells from multiple columns in an xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        columns = config.get('columns', ['B', 'C'])
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 8)
        cells = {}
        for col in columns:
            for row in range(start_row, end_row + 1):
                cell_ref = f'{col}{row}'
                val = ws[cell_ref].value
                cells[cell_ref] = str(val) if val is not None else None
        return {'cells': cells}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_header_and_values__23bfc830003e8f489b53a39808f56b59(env, config: dict):
    """Get header and selected cell values from a column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        col_letter = config.get('column', 'E')
        header_row = config.get('header_row', 1)
        check_rows = config.get('check_rows', [2, 10, 20, 30])
        header = ws[f'{col_letter}{header_row}'].value
        values = {}
        for r in check_rows:
            cell_val = ws[f'{col_letter}{r}'].value
            values[str(r)] = cell_val
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3440d9b31b1654ffe8da4da80be4e5fc(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__52b5a1c85fa3a63d940f3aeb966674a8(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'D1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3d7517c3edf35996bcb8466dea369ae2(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__8ed2a3a912f306294d5ece5e575ff288(env, config: dict):
    """Get value of a specific cell from xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'C1')
        value = ws[cell].value
        if value is not None:
            try:
                value = float(value)
            except (ValueError, TypeError):
                pass
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cumulative_col__6bc8ae788974152807b9723497a26da7(env, config: dict):
    """Get column D header and values for cumulative sales check."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        header = ws['D1'].value
        values = []
        for row in range(2, 12):
            val = ws.cell(row=row, column=4).value
            values.append(val)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__34bc68338cfc2efa5ed7f80b3fe4afc3(env, config: dict):
    """Get all values from a specific column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        col_letter = config.get('column', 'F')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws[f'{col_letter}{row}'].value
            if val is not None:
                values.append(val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cells__677226b05224be2ecbd2de2e693f9c8e(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        if isinstance(sheet_idx, int):
            ws = wb.worksheets[sheet_idx]
        else:
            ws = wb[sheet_idx]
        cells = config.get('cells', [])
        values = {}
        for cell in cells:
            val = ws[cell].value
            values[cell] = val
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__5eff620ab5e8861fd0f5d3b257b5ca40(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'E1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_values__e0472552b507a6465c494b007b9b0305(env, config: dict):
    """Get values from specific cells in a row range from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cells = config.get('cells', [])
        values = {}
        for cell_ref in cells:
            val = ws[cell_ref].value
            values[cell_ref] = val
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__2b83d943849f7c327e47bb940fe75967(env, config: dict):
    """Get a range of cell values from a column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        column = config.get('column', 'E')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 7)
        values = []
        for row in range(start_row, end_row + 1):
            cell_ref = f'{column}{row}'
            val = ws[cell_ref].value
            values.append(val)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sum_row__99d61f333e953d709c5fe5f221bfb63e(env, config: dict):
    """Get values from the total/sum row in the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 12)
        label = ws.cell(row=row, column=1).value
        sales_total = ws.cell(row=row, column=2).value
        cogs_total = ws.cell(row=row, column=3).value
        return {'label': label, 'sales_total': sales_total, 'cogs_total': cogs_total}
    finally:
        os.unlink(tmp_path)

def get_timetable_cell__929c2cac5efeaad2cf5e1556b5f1f38f(env, config: dict):
    """Get cell value and background color from the Course Timetable xlsx."""
    file_path = config.get('path', '/home/user/Desktop/Course Timetable.xlsx')
    cell_addr = config.get('cell', 'C6')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell = ws[cell_addr]
        value = cell.value
        fill = cell.fill
        bg_color = 'none'
        if fill and fill.start_color and fill.start_color.rgb:
            rgb = fill.start_color.rgb
            if rgb and rgb != '00000000':
                bg_color = str(rgb)
        return {'value': value, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_xlsx_names_and_city__95be0722e22941aff21eff3d680c0c01(env, config: dict):
    """Get restaurant names (A2:A6) and city column (E1:E6) from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        names = []
        for row in range(2, 7):
            val = ws.cell(row=row, column=1).value
            names.append(str(val).strip() if val else None)
        e1_val = ws['E1'].value
        city_values = []
        for row in range(2, 7):
            val = ws.cell(row=row, column=5).value
            city_values.append(str(val).strip() if val else None)
        return {'names': names, 'e1_header': str(e1_val).strip() if e1_val else None, 'city_values': city_values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__da4a3fddcbcd3738d01b04b1ba353fc4(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__56ea06753821989747521678ff3594b3(env, config: dict):
    """Get all values from a specific column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 3)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            cell_val = ws.cell(row=row, column=col).value
            values.append(cell_val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells__bf5b353ae5739d8c5c981c28dcf67581(env, config: dict):
    """Get values from specific cells in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        result = {}
        for cell_ref in config.get('cells', ['A9', 'D9']):
            result[cell_ref] = ws[cell_ref].value
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3a7517792dc1ec4e3be4f857bdb0950f(env, config: dict):
    """Get a cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_data__0d0f959234d0c494b6d7c17cb5158836(env, config: dict):
    """Get all data from a specific sheet in an xlsx file."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', None)
        if sheet_name and sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheets': wb.sheetnames}
        ws = wb[sheet_name] if sheet_name else wb.worksheets[0]
        headers = []
        for col in range(1, ws.max_column + 1):
            headers.append(ws.cell(1, col).value)
        rows = []
        for row in range(2, ws.max_row + 1):
            row_data = {}
            for col in range(1, ws.max_column + 1):
                row_data[headers[col - 1]] = ws.cell(row, col).value
            if any((v is not None for v in row_data.values())):
                rows.append(row_data)
        return {'sheet_name': ws.title, 'headers': headers, 'rows': rows, 'row_count': len(rows), 'all_sheets': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_zone2_col_totals__dfa00147f689f96c837b3a22754c016b(env, config: dict):
    """Get the Total row (row 13) values for Zone 2 quarterly columns (B-E)."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for col in ['B', 'C', 'D', 'E']:
            cell_ref = f'{col}13'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__3dfc90e20421740ec2b3449751e1ede4(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sorted_sales__38a9aeefb36e58a6a0dcece70b9af3ed(env, config: dict):
    """Get Sales column values after sorting to verify sort order."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        sales_values = []
        week_labels = []
        for row in range(2, 12):
            week_labels.append(ws.cell(row=row, column=1).value)
            sales_values.append(ws.cell(row=row, column=2).value)
        return {'week_labels': week_labels, 'sales_values': sales_values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_clevel_count__802aa06dfc0f94fbf79660d5d540798f(env, config: dict):
    """Get the C-level count cell (E2) and the rank column (D) from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {'e1_value': ws.cell(row=1, column=5).value, 'e2_value': ws.cell(row=2, column=5).value, 'ranks': []}
        for row in range(2, 23):
            result['ranks'].append(ws.cell(row=row, column=4).value)
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__27fb6350ba8922eab313e55ef2751ee9(env, config: dict):
    """Get a specific cell value from an xlsx file on the VM."""
    file_path = config.get('path', '')
    cell = config.get('cell', 'A1')
    sheet_index = config.get('sheet', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_index]
        value = ws[cell].value
        return {'value': value, 'cell': cell}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__0762adbae4ad235c185f47096cf64c91(env, config: dict):
    """Get a range of cell values from a column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        col = config.get('column', 5)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            v = ws.cell(row=row, column=col).value
            if v is not None:
                values.append(str(v).strip())
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__40424b53798c2f42a45c9b5d1f78c9a0(env, config: dict):
    """Get a range of cell values from a column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        column = config.get('column', 7)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=column).value
            values.append(val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_sales_by_rep__5c7590b2fce2859cb43c07916cdde25c(env, config: dict):
    """Read Sheet2 to get sales rep names and their total sales values."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if len(wb.sheetnames) < 2:
            return {'error': 'Sheet2 not found'}
        ws = wb.worksheets[1]
        data = {}
        for row in range(1, ws.max_row + 1):
            for col in range(1, ws.max_column + 1):
                val = ws.cell(row=row, column=col).value
                if val is not None:
                    val_str = str(val).strip()
                    if val_str in ('Joe', 'Moe', 'Chin'):
                        for c2 in range(1, ws.max_column + 1):
                            if c2 != col:
                                v2 = ws.cell(row=row, column=c2).value
                                if isinstance(v2, (int, float)) and v2 > 0:
                                    data[val_str] = float(v2)
                                    break
        return {'sales_by_rep': data}
    finally:
        os.unlink(tmp_path)

def get_timetable_cell__ae3105aee7e1d6920d1ca91156ffa145(env, config: dict):
    """Get cell value and background color from the Course Timetable xlsx."""
    file_path = config.get('path', '/home/user/Desktop/Course Timetable.xlsx')
    cell_addr = config.get('cell', 'E3')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell = ws[cell_addr]
        value = cell.value
        fill = cell.fill
        bg_color = 'none'
        if fill and fill.start_color and fill.start_color.rgb:
            rgb = fill.start_color.rgb
            if rgb and rgb != '00000000':
                bg_color = str(rgb)
        return {'value': value, 'bg_color': bg_color}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__9e8bcb68d0493b0f9f821969711fa470(env, config: dict):
    """Get a single cell value from an Excel file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cell = config.get('cell', 'B11')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__8afea61a57c3ecda7ab6098e540526c4(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_path = config.get('path', '/home/user/Desktop/researchers.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_reversed_names__d8dd1c41700e60d858e6fd9c2f1f5451(env, config: dict):
    """Get column E (reversed full name) and columns B, C from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {'header_e': ws.cell(row=1, column=5).value, 'full_names': [], 'first_names': [], 'last_names': []}
        for row in range(2, 23):
            result['full_names'].append(ws.cell(row=row, column=5).value)
            result['first_names'].append(ws.cell(row=row, column=2).value)
            result['last_names'].append(ws.cell(row=row, column=3).value)
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_d_values__d026dfbc396400238630feee90a0bfc7(env, config: dict):
    """Read values from column D (Pass/Fail/Held) for rows 2-29."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        values = {}
        for row in range(2, 30):
            cell_val = ws.cell(row=row, column=4).value
            values[f'D{row}'] = cell_val
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__02d2303f1946e73506bf1a1871c3326e(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__7b9538bff9f96ac19c0aebd38f8a5f3f(env, config: dict):
    """Get sheet names from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_maxmin__56726039aefd1f0c57475732eb3ab1f0(env, config: dict):
    """Get max Revenue and min Expenses from Sheet2."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet2 = None
        for ws in wb.worksheets:
            if 'sheet2' in ws.title.lower().replace(' ', ''):
                sheet2 = ws
                break
        if sheet2 is None and len(wb.worksheets) > 1:
            sheet2 = wb.worksheets[1]
        if sheet2 is None:
            return {'error': 'Sheet2 not found'}
        header_a = sheet2.cell(row=1, column=1).value
        header_b = sheet2.cell(row=1, column=2).value
        value_a = sheet2.cell(row=2, column=1).value
        value_b = sheet2.cell(row=2, column=2).value
        return {'header_a': header_a, 'header_b': header_b, 'value_a': value_a, 'value_b': value_b}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_cell__1546f90c5e4d91a5767b01e5f0a56119(env, config: dict):
    """Get text from a specific table cell in a pptx slide."""
    from pptx import Presentation
    file_path = config.get('path', '')
    slide_index = config.get('slide_index', 0)
    row_index = config.get('row_index', 0)
    col_index = config.get('col_index', 0)
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        prs = Presentation(tmp_path)
        if slide_index >= len(prs.slides):
            return {'error': f'Slide index {slide_index} out of range'}
        slide = prs.slides[slide_index]
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                if row_index < len(table.rows) and col_index < len(table.columns):
                    cell_text = table.cell(row_index, col_index).text.strip()
                    return {'cell_text': cell_text}
                else:
                    return {'error': f'Cell [{row_index},{col_index}] out of range'}
        return {'error': 'No table found on slide'}
    finally:
        os.unlink(tmp_path)

def get_xlsx_freeze_pane__7c0a43d93f3ca3bc0e18377c9b4b22b9(env, config: dict):
    """Get the freeze pane setting of the specified sheet in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        freeze_panes = ws.freeze_panes
        return {'freeze_panes': str(freeze_panes) if freeze_panes else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__542fc695553e3e27095b27327032df35(env, config: dict):
    """Get all values from a specified column in an xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 1)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(str(val) if val is not None else None)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_data__c2e2ad1d381d1d27297ba9392b0d9cea(env, config: dict):
    """Get all data from a specific sheet in an xlsx file."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', None)
        if sheet_name and sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheets': wb.sheetnames}
        ws = wb[sheet_name] if sheet_name else wb.worksheets[0]
        headers = []
        for col in range(1, ws.max_column + 1):
            headers.append(ws.cell(1, col).value)
        rows = []
        for row in range(2, ws.max_row + 1):
            row_data = {}
            for col in range(1, ws.max_column + 1):
                key = str(headers[col - 1]).strip().lower() if headers[col - 1] is not None else ''
                row_data[key] = ws.cell(row, col).value
            if any((v is not None for v in row_data.values())):
                rows.append(row_data)
        return {'sheet_name': ws.title, 'headers': headers, 'rows': rows, 'row_count': len(rows), 'all_sheets': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__427c1b1ebceb5c2ead27edb218ee7bf9(env, config: dict):
    """Get values from a column range in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        column = config.get('column', 'F')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 23)
        values = []
        for row in range(start_row, end_row + 1):
            cell_ref = f'{column}{row}'
            values.append(ws[cell_ref].value)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__12cf17c32a6915e1266996d25d563e2a(env, config: dict):
    """Get a cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_values__8c5594571c596e8d15ba168f23066d41(env, config: dict):
    """Get values from a specific row in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 12)
        cols = config.get('cols', ['A', 'B', 'C', 'D', 'E', 'F', 'G'])
        result = {}
        for col in cols:
            cell_ref = f'{col}{row}'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_top_performers__85190e1339a2cb1782e74f01270bc172(env, config: dict):
    """Get Month and TopRep columns from TopPerformers sheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        target_sheet = None
        for name in wb.sheetnames:
            if 'top' in name.lower() and 'perform' in name.lower():
                target_sheet = wb[name]
                break
        if target_sheet is None:
            return {'error': 'TopPerformers sheet not found'}
        month_top = {}
        for row in range(2, target_sheet.max_row + 1):
            month = target_sheet.cell(row=row, column=1).value
            rep = target_sheet.cell(row=row, column=2).value
            if month is None:
                break
            month_top[str(month).strip()] = str(rep).strip() if rep else None
        headers = []
        for col in range(1, target_sheet.max_column + 1):
            val = target_sheet.cell(row=1, column=col).value
            if val is not None:
                headers.append(str(val).strip().lower())
        return {'sheet_found': True, 'headers': headers, 'month_top': month_top, 'row_count': len(month_top)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_data__00c370b29bbdeefb624743c61a277c9e(env, config: dict):
    """Get cell values from a specific row in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 9)
        result = {}
        for col_letter in config.get('columns', ['A', 'B', 'C', 'D', 'E']):
            cell = ws[f'{col_letter}{row}']
            result[col_letter] = cell.value
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_values__767f7417cb0fc080051f9d7103e674cb(env, config: dict):
    """Get values from a specific row in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        row = config.get('row', 12)
        cols = config.get('cols', ['A', 'B', 'C', 'D', 'E', 'F', 'G'])
        result = {}
        for col in cols:
            cell_ref = f'{col}{row}'
            val = ws[cell_ref].value
            result[cell_ref] = val
        return result
    finally:
        os.unlink(tmp_path)

def get_ods_to_xlsx_check__d8b4b0a7ce586e5e25f592627cc186b5(env, config: dict):
    """Check if ODS was converted to XLSX via terminal command."""
    result = {}
    try:
        history_output = env.controller.run_bash_script("cat ~/.bash_history | grep '\\(soffice\\|libreoffice\\).*--convert-to\\s\\+xlsx'", timeout=30)
        history_text = history_output.get('output', '') if isinstance(history_output, dict) else str(history_output)
        result['used_terminal'] = 'use terminal' if history_text.strip() else 'use no terminal'
    except Exception:
        result['used_terminal'] = 'use no terminal'
    xlsx_path = config.get('path', '/home/user/Desktop/file_example_ODS_5000.xlsx')
    try:
        file_bytes = env.controller.get_file(xlsx_path)
        if file_bytes and len(file_bytes) > 0:
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                import openpyxl
                wb = openpyxl.load_workbook(tmp_path, data_only=True)
                ws = wb.active
                result['file_exists'] = True
                result['row_count'] = ws.max_row
                result['col_count'] = ws.max_column
                first_data_cell = ws.cell(row=2, column=2).value
                result['has_data'] = first_data_cell is not None
            finally:
                os.unlink(tmp_path)
        else:
            result['file_exists'] = False
            result['row_count'] = 0
            result['col_count'] = 0
            result['has_data'] = False
    except Exception:
        result['file_exists'] = False
        result['row_count'] = 0
        result['col_count'] = 0
        result['has_data'] = False
    return result

def get_xlsx_cell_value__4eed649bcee7f8431b3567000e60f20f(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'D1')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_col_e_values__c2b3326804f7e34493f7827960c36c6d(env, config: dict):
    """Read column E values (rows 2-30) from the spreadsheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws.cell(row=1, column=5).value
        values = []
        for row in range(2, 31):
            cell = ws.cell(row=row, column=5)
            val = cell.value
            if val is not None:
                if isinstance(val, float) and val == int(val):
                    values.append(int(val))
                else:
                    values.append(val)
            else:
                values.append(None)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cs_data__6d4ad198dfe8ff0251a084115fc1444c(env, config: dict):
    """Get CS application data from xlsx file."""
    import openpyxl
    file_path = config.get('path', '/home/user/Desktop/GRF-cs-2023.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        header_a = ws.cell(row=1, column=1).value
        header_b = ws.cell(row=1, column=2).value
        value_a = ws.cell(row=2, column=1).value
        value_b = ws.cell(row=2, column=2).value
        return {'header_a': str(header_a).strip() if header_a is not None else None, 'header_b': str(header_b).strip() if header_b is not None else None, 'value_a': value_a, 'value_b': value_b}
    finally:
        os.unlink(tmp_path)

def get_xlsx_two_cells__009b0a2429b483dc78b74a0168744c1a(env, config: dict):
    """Get values from two specified cells in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell1 = config.get('cell1', 'F2')
        cell2 = config.get('cell2', 'G2')
        return {'name': ws[cell1].value, 'count': ws[cell2].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__0f45568bd733d9527d872ad7c053c2c8(env, config: dict):
    """Get all values from a specified column in an xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 2)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__67a1f86b27a24adb60e8a183c35b0852(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_path = config.get('path', '/home/user/DemographicProfile.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet', 'Sheet1')
        if sheet_name not in wb.sheetnames:
            return {'error': f"Sheet '{sheet_name}' not found"}
        ws = wb[sheet_name]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__4abf9f9f6bbbc1d8a7c4b9c7e898d79c(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sort_check__08b8e71e292866610cba11baa24dc9d9(env, config: dict):
    """Get first several cell values from column A to verify sort order."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        dates = []
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=1).value
            if val is not None:
                dates.append(str(val))
            else:
                dates.append(None)
        is_sorted = True
        for i in range(len(dates) - 1):
            if dates[i] is not None and dates[i + 1] is not None:
                if dates[i] > dates[i + 1]:
                    is_sorted = False
                    break
        first_val = dates[0] if dates else None
        last_val = dates[-1] if dates else None
        return {'is_sorted_ascending': is_sorted, 'first_value': first_val, 'last_value': last_val, 'row_count': len(dates)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__990c97f8e62011f017525e7c1f376ba9(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__ad5647b3d1b47a9e1996a6e69ca9562d(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__28ae1ff6307904e177d5ad2de88ea142(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__69f94a2f4d67df544b1a5f49c4c42c6d(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__9889ac58b2fb5d2c955b2821b86f86b7(env, config: dict):
    """Get sheet names from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_number_format__092734e9cc221c6f71e099e209498fb4(env, config: dict):
    """Get the number format of specified cells from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells_to_check = config.get('cells', [])
        formats = {}
        for cell_ref in cells_to_check:
            cell = ws[cell_ref]
            formats[cell_ref] = cell.number_format
        return {'formats': formats}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_bgcolor__0e55cf8187f0e07889b1109cdd0b266f(env, config: dict):
    """Get background colors of specified cells from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_index = config.get('sheet', 0)
        ws = wb.worksheets[sheet_index]
        cells = config.get('cells', [])
        colors = {}
        for cell_ref in cells:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and fill.fgColor.rgb and (fill.fgColor.rgb != '00000000'):
                colors[cell_ref] = str(fill.fgColor.rgb)
            elif fill and fill.patternType == 'solid' and fill.start_color and fill.start_color.rgb:
                colors[cell_ref] = str(fill.start_color.rgb)
            else:
                colors[cell_ref] = None
        return {'colors': colors, 'count': len(colors)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__6d702aff56732f3b1cb6a80e7f740e70(env, config: dict):
    """Get all values from a specific column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 3)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', ws.max_row)
        values = []
        for row in range(start_row, end_row + 1):
            cell_val = ws.cell(row=row, column=col).value
            values.append(cell_val)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__2650b8c230f7272f8553cf2f429cf59a(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__9ff5dc914928ae745b032698ddd720c0(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_b3__0698f7f5da73d25e009ccfe66a594e11(env, config: dict):
    """Get the value of cell B3 from the spreadsheet."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws['B3'].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__2fc57009f3993d012dcfddb9200ab322(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_merge_text__221a8bb1dbb77d2fdb6346ea5a7576e4(env, config: dict):
    """Get Sheet2 status: existence, merged cells, and A1 value."""
    file_path = config.get('path', '/home/user/DemographicProfile.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'sheet2_exists': False, 'a1_value': None, 'is_merged': False}
        ws = wb['Sheet2']
        a1_value = ws['A1'].value
        is_merged = False
        for merged_range in ws.merged_cells.ranges:
            if str(merged_range) == 'A1:C1':
                is_merged = True
                break
        return {'sheet2_exists': True, 'a1_value': a1_value, 'is_merged': is_merged}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_values__3d93b5527badb1dca77f44852a74fbe3(env, config: dict):
    """Read values from specified cells."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cells = config.get('cells', [])
        values = {}
        for cell_ref in cells:
            values[cell_ref] = ws[cell_ref].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__949a54a656a59ff2d23607cb949dacfc(env, config: dict):
    """Get a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'B8')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell__8a05426d089eda14388178863fc78e0b(env, config: dict):
    """Get multiple cell values from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cells = config.get('cells', [])
        result = {}
        for cell_ref in cells:
            result[cell_ref] = ws[cell_ref].value
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__0ff3d8d538e7cd587b08458b9c790695(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sorted_names__56c311ddff224f797b6614d5151a3c35(env, config: dict):
    """Get employee names in column A after sorting to verify alphabetical order."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        names = []
        for row in range(2, 9):
            val = ws.cell(row=row, column=1).value
            names.append(val)
        return {'names': names}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__cb5b990d93f87d0f04a6150ee1d23a05(env, config: dict):
    """Get a specific cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_range_values__160ec07bd1e641be9b31cf92f3848b19(env, config: dict):
    """Get values from a range of cells in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cells = config.get('cells', [])
        values = {}
        for cell in cells:
            values[cell] = ws[cell].value
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__d5a8ae815c3da82f8081ee65f4e2cea1(env, config: dict):
    """Read a single cell value from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'D1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_all_data__d9fdc5409ef8a872907b1ddc3db9cd54(env, config: dict):
    """Read all data rows from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        rows = []
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column, values_only=True):
            rows.append(list(row))
        return {'rows': rows, 'max_row': ws.max_row}
    finally:
        os.unlink(tmp_path)

def get_xlsx_number_format__6f34d4ab7744f84c76fe24f84ac55b09(env, config: dict):
    """Get number format of cells in a column."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=False)
        ws = wb.worksheets[config.get('sheet', 0)]
        col_letter = config.get('column', 'C')
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 8)
        formats = []
        values = []
        for row in range(start_row, end_row + 1):
            cell = ws[f'{col_letter}{row}']
            formats.append(cell.number_format)
            values.append(cell.value)
        return {'formats': formats, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_half_rate_col__fd58351f5fafe47d305d3595c5202dcd(env, config: dict):
    """Get column C header and half-rate values from PeriodRate.xlsx."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws['C1'].value
        values = {}
        for row in range(2, 26):
            cell_val = ws.cell(row=row, column=3).value
            if cell_val is not None:
                values[f'C{row}'] = round(float(cell_val), 4)
            else:
                values[f'C{row}'] = None
        return {'header': header, 'values': values, 'count': sum((1 for v in values.values() if v is not None))}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet2_avg__843226a2563502226c076bf6effbabf5(env, config: dict):
    """Get average values from Sheet2."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet2 = None
        for ws in wb.worksheets:
            if 'sheet2' in ws.title.lower().replace(' ', ''):
                sheet2 = ws
                break
        if sheet2 is None and len(wb.worksheets) > 1:
            sheet2 = wb.worksheets[1]
        if sheet2 is None:
            return {'error': 'Sheet2 not found'}
        header_a = sheet2.cell(row=1, column=1).value
        header_b = sheet2.cell(row=1, column=2).value
        value_a = sheet2.cell(row=2, column=1).value
        value_b = sheet2.cell(row=2, column=2).value
        return {'header_a': header_a, 'header_b': header_b, 'value_a': value_a, 'value_b': value_b}
    finally:
        os.unlink(tmp_path)

def get_xlsx_net_income_column__bc1bb048c0eeb30a1fa5f9604d4f0a11(env, config: dict):
    """Get column C (Net Income) header and values from Sheet1."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws.cell(row=1, column=3).value
        values = []
        for row in range(2, 21):
            val = ws.cell(row=row, column=3).value
            values.append(val)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sort_order__db00ed878507fde13aefe6b71212bcbd(env, config: dict):
    """Get the values in column B (All citations) to check sort order."""
    file_path = config.get('path', '/home/user/Desktop/researchers.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        rows = []
        for row in range(2, ws.max_row + 1):
            name = ws.cell(row=row, column=1).value
            citations = ws.cell(row=row, column=2).value
            if name is not None and citations is not None:
                rows.append({'name': name, 'citations': citations})
        return {'rows': rows, 'count': len(rows)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_by_label__4551fb3f098ec9fe882ac83568cd431f(env, config: dict):
    """Scan column A for a label and return values from columns B-G of that row."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        label = config.get('label', 'Total')
        target_row = None
        for row in range(1, ws.max_row + 1):
            cell_val = ws.cell(row=row, column=1).value
            if cell_val and str(cell_val).strip().lower() == label.lower():
                target_row = row
                break
        if target_row is None:
            return {'error': f'Label "{label}" not found in column A', 'label_found': False}
        values = {}
        col_names = ['B', 'C', 'D', 'E', 'F', 'G']
        for (i, col_name) in enumerate(col_names):
            val = ws.cell(row=target_row, column=i + 2).value
            values[col_name] = val
        return {'label_found': True, 'label': str(ws.cell(row=target_row, column=1).value).strip(), 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__54e9a68f9afb8559eff309088793c132(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_values__70ea2799be65bb426ef2e5f3f76ece43(env, config: dict):
    """Read all values from a specified column in an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        col = config.get('column', 3)
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 22)
        values = []
        for row in range(start_row, end_row + 1):
            val = ws.cell(row=row, column=col).value
            values.append(str(val).strip() if val is not None else None)
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__fec4261d3cbd8743e0d3445b3fa0111f(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        if isinstance(sheet_idx, int):
            ws = wb.worksheets[sheet_idx]
        else:
            ws = wb[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_qty_sort_check__d422460342c18f4ab4e35ab11ffce7d8(env, config: dict):
    """Get column E (Quantity) values to verify descending sort order."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        quantities = []
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=5).value
            if val is not None:
                quantities.append(val)
        is_sorted_desc = True
        for i in range(len(quantities) - 1):
            if quantities[i] < quantities[i + 1]:
                is_sorted_desc = False
                break
        first_qty = quantities[0] if quantities else None
        row_count = len(quantities)
        return {'is_sorted_descending': is_sorted_desc, 'first_quantity': first_qty, 'row_count': row_count}
    finally:
        os.unlink(tmp_path)

def get_xlsx_col_c_values__85a76f9ca092376f5f612fc019e6af39(env, config: dict):
    """Read column C (Old ID) values from rows 2-30 to check sort order."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        values = []
        for row in range(2, 31):
            cell = ws.cell(row=row, column=3)
            val = cell.value
            if val is not None:
                if isinstance(val, float) and val == int(val):
                    values.append(int(val))
                else:
                    values.append(val)
            else:
                values.append(None)
        values = [v for v in values if v is not None]
        return {'values': values, 'count': len(values)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__d3f5a4ed5eae2c51d540e2a10864dc43(env, config: dict):
    """Get the value of a specific cell from an xlsx file."""
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__ab8a45e026223c32b7b22c741f1bc3ca(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        if isinstance(sheet_idx, int):
            ws = wb.worksheets[sheet_idx]
        else:
            ws = wb[sheet_idx]
        cell = config.get('cell', 'A1')
        value = ws[cell].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_fv_column__5c1e2dfe8bffbdb20d88aa7499677c46(env, config: dict):
    """Get column F values from Sheet1 (Future Value column)."""
    file_path = config.get('path', '/home/user/FutureValue.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        result['F1'] = ws['F1'].value
        for row in range(2, 6):
            cell = ws.cell(row=row, column=6)
            result[f'F{row}'] = cell.value
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__e7bff0bae84b8431a03877cfc4fe4751(env, config: dict):
    """Read specific cells from an xlsx file."""
    import openpyxl
    path = config.get('path', '')
    cells = config.get('cells', ['C7'])
    sheet_index = config.get('sheet', 0)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_index]
        result = {}
        for cell_ref in cells:
            result[cell_ref] = ws[cell_ref].value
        total_rows = []
        for row_idx in range(1, ws.max_row + 1):
            a_val = ws.cell(row=row_idx, column=1).value
            c_val = ws.cell(row=row_idx, column=3).value
            if a_val and 'total' in str(a_val).lower():
                total_rows.append({'row': row_idx, 'label': a_val, 'amount': c_val})
        result['total_rows'] = total_rows
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_column_cells__4065e6a4b9dbccc0390b8c5a38e2d1b7(env, config: dict):
    """Read a range of cells from a specific column in an xlsx file."""
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        column = config.get('column', 'D')
        start_row = config.get('start_row', 1)
        end_row = config.get('end_row', ws.max_row)
        cells = {}
        for row in range(start_row, end_row + 1):
            cell_ref = f'{column}{row}'
            val = ws[cell_ref].value
            cells[cell_ref] = str(val) if val is not None else None
        return {'cells': cells}
    finally:
        os.unlink(tmp_path)

def get_sheet_names__a164952ae6b41142a59183faf6bedadf_qw35sft2_0b7b29d8(env, config: dict):
    """Get the ordered list of sheet names from copy_sheet_insert.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/copy_sheet_insert.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_merged_cell_state__48a88b5bb3b362cc9c78cc671cedfc7c_qw35sft2_80edd89c(env, config: dict):
    """Get merged cell info and cell value from a named sheet in an xlsx file."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', '')
        if sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheet_names': wb.sheetnames}
        ws = wb[sheet_name]
        merge_range = config.get('merge_range', '')
        merged_ranges = [str(m) for m in ws.merged_cells.ranges]
        is_merged = merge_range in merged_ranges if merge_range else False
        cell_ref = config.get('cell', 'A1')
        value = ws[cell_ref].value
        return {'is_merged': is_merged, 'value': value, 'merged_ranges': merged_ranges}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__0762adbae4ad235c185f47096cf64c91_qw35sft2_84cd3d68(env, config: dict):
    """Read value of a specific cell from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Names_Duplicate_Unique.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell_ref = config.get('cell', 'D2')
        return {'value': ws[cell_ref].value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_state__68824c07a6606ac16add37bf4765401b_qw35sft2_ca181f23(env, config: dict):
    """Get Sheet2 existence and cell A1 value from DemographicProfile.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/DemographicProfile.xlsx')
    if not file_bytes:
        return {'error': 'File not found', 'sheet_exists': False, 'a1_value': None}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_exists = 'Sheet2' in wb.sheetnames
        a1_value = None
        if sheet_exists:
            ws2 = wb['Sheet2']
            a1_value = ws2['A1'].value
        return {'sheet_exists': sheet_exists, 'a1_value': a1_value}
    finally:
        os.unlink(tmp_path)

def get_seqno_and_sheet_name__83227706efd58da00bb8213cee393c7a_qw35sft2_a23bfda1(env, config: dict):
    """Read Seq No. column (B2:B29) and the first sheet's name from the xlsx file."""
    path = config.get('path', '/home/user/Order_Sales_Serial#.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        seq_nos = [ws.cell(row=r, column=2).value for r in range(2, 30)]
        sheet_name = ws.title
        return {'seq_nos': seq_nos, 'sheet_name': sheet_name}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__cfba55e4566ed0373f9b47631b661201_qw35sft2_b5a6828b(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        return {'value': ws[config.get('cell', 'A1')].value}
    finally:
        os.unlink(tmp_path)

def get_calc_len_formula__29aa3f16a2f40add619bb0073f57ed2e_qw35sft2_2fc070d1(env, config: dict):
    """Read C2 and D2 values from the movie titles spreadsheet."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Movie_title_garbage_clean.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        c2_value = ws['C2'].value
        d2_value = ws['D2'].value
        return {'c2_value': c2_value, 'd2_value': d2_value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_level_secondary__f558bfbb444b8f5736dbe2588385d5fe_qw35sft2_c0e794c7(env, config: dict):
    """Read B8:B18 cell values from Student_Level_Fill_Blank.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/Student_Level_Fill_Blank.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {f'B{r}': ws[f'B{r}'].value for r in range(8, 19)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__d4b207925bd2d0c3ababe1e319fa8b1f_qw35sft2_70ffb6dd(env, config: dict):
    """Read a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Create_column_charts_using_statistics.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'value': None}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell_ref = config.get('cell', 'K2')
        value = ws[cell_ref].value
        return {'value': value, 'cell': cell_ref}
    except Exception as e:
        return {'error': str(e), 'value': None}
    finally:
        os.unlink(tmp_path)

def get_calc_cell_e3__73e90374b4fd7034756570a58d380c35_qw35sft2_d5812f2b(env, config: dict):
    """Read numeric value of cell E3 from the saved xlsx file."""
    path = config.get('path', '/home/user/Multiply_Time_Number.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        value = ws['E3'].value
        return {'e3_value': float(value) if value is not None else None}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__1ed2b6f96cb9e3177c42c0765213bb8c_qw35sft2_a3c8bc99(env, config: dict):
    """Read a specific cell value from an xlsx file on the VM."""
    path = config.get('path', '/home/user/Date_Budget_Variance_HideNA.xlsx')
    cell_ref = config.get('cell', 'A1')
    sheet_idx = config.get('sheet', 0)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[sheet_idx]
        value = ws[cell_ref].value
        return {'value': value, 'cell': cell_ref}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cell__3b5b28cd38bc58813bd97e9b2f5d4ad4_qw35sft2_ebc90768(env, config: dict):
    """Read multiple cell values from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cells = config.get('cells', [])
        return {cell: ws[cell].value for cell in cells}
    finally:
        os.unlink(tmp_path)

def get_xlsx_avg_age_cell__c483fe0106f1e0b2430ad19cd97c77c0_qw35sft2_d43a6ec6(env, config: dict):
    """Read cell E3 from Zoom_Out_Oversized_Cells.xlsx on VM (expected to hold average age formula)."""
    file_bytes = env.controller.get_file('/home/user/Zoom_Out_Oversized_Cells.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'E3': ws['E3'].value}
    finally:
        os.unlink(tmp_path)

def get_calc_sheet1_net_income__6142d421fefd6b784b8ad81070a58350_qw35sft2_1a6fd6f4(env, config: dict):
    """Read Sheet1 column C header and net income values (rows 2-20)."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/NetIncome.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet1' not in wb.sheetnames:
            return {'error': 'Sheet1 not found'}
        ws = wb['Sheet1']
        header = ws['C1'].value
        values = []
        for row in range(2, 21):
            v = ws.cell(row=row, column=3).value
            if isinstance(v, float) and v == int(v):
                v = int(v)
            values.append(v)
        return {'header': header, 'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sort_total_row__2cdc28af4ef71e22984f659197334a50_qw35sft2_15e48de1(env, config: dict):
    """Read D2, A20, and D20 to verify sort + Total row was added correctly."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Arrang_Value_min_to_max.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'d2': ws.cell(row=2, column=4).value, 'a20': ws.cell(row=20, column=1).value, 'd20': ws.cell(row=20, column=4).value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_sorted_revenue__ff6f4e09739d1959467ff60470ffe2bd_qw35sft2_69d06ae7(env, config: dict):
    """Get Sheet2 Revenue column header, row count, sort order, and boundary values."""
    file_bytes = env.controller.get_file('/home/user/NetIncome.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': list(wb.sheetnames)}
        ws = wb['Sheet2']
        header = ws['A1'].value
        data = [ws.cell(row=r, column=1).value for r in range(2, 21)]
        non_null = [v for v in data if v is not None]
        is_sorted_asc = all((non_null[i] <= non_null[i + 1] for i in range(len(non_null) - 1))) if len(non_null) > 1 else True
        return {'header': header, 'row_count': len(non_null), 'is_sorted_asc': is_sorted_asc, 'first_value': non_null[0] if non_null else None, 'last_value': non_null[-1] if non_null else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_b7__9ab5f66fb165e5668bd5b38cac0c73c2_qw35sft2_bb64f4bd(env, config: dict):
    """Read cell B7 value from Represent_in_millions_billions.xlsx."""
    import tempfile
    import os
    import openpyxl
    file_path = config.get('path', '/home/user/Represent_in_millions_billions.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        value = ws['B7'].value
        return {'value': str(value) if value is not None else ''}
    finally:
        os.unlink(tmp_path)

def get_xlsx_first_col_header__b3e437bf5bf6ec56d9f1c4d46601f3ff_qw35sft2_64bd5de5(env, config: dict):
    """Get the header value in cell A1 of the spreadsheet."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Name_Order_Id_move_column.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        a1_value = ws.cell(row=1, column=1).value
        return {'a1_header': a1_value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_colors__03a99dd86c6d30983d2467bc1177489c_qw35sft2_df3a9ab2(env, config: dict):
    """Get background fill colors for specified cells in an xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Calendar_Highlight_Weekend_Days.xlsx')
    cells_to_check = config.get('cells', [])
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for cell_ref in cells_to_check:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and (fill.fgColor.type == 'rgb'):
                color = fill.fgColor.rgb
            else:
                color = '00000000'
            result[cell_ref] = color
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_transposed_cells__f7593e4503a45a853c1ae96cb08aaf92_qw35sft2_3f7147c6(env, config: dict):
    """Read key cells from the transposed table region B8:E12 in Students_Class_Subject_Marks.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Students_Class_Subject_Marks.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'B8': ws['B8'].value, 'E8': ws['E8'].value, 'B9': ws['B9'].value, 'E9': ws['E9'].value, 'B12': ws['B12'].value, 'E12': ws['E12'].value}
    finally:
        os.unlink(tmp_path)

def get_calc_sheet_info__ecf996871673b767fb6d252021f29e29_qw35sft2_7e7fd42c(env, config: dict):
    """Read sheet names from WeeklySales.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/WeeklySales.xlsx')
    if not file_bytes:
        return {'error': 'File not found', 'sheet_names': []}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': list(wb.sheetnames)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_zone2_row_totals__6ef64aee01b9e4fef352e390ba82e0ce_qw35sft2_14275934(env, config: dict):
    """Read Zone 2 product row totals (F10, F11, F12) from Quarterly_Product_Sales_by_Zone.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Quarterly_Product_Sales_by_Zone.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'F10': ws['F10'].value, 'F11': ws['F11'].value, 'F12': ws['F12'].value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_three_columns__74fb6bc2ea3e8707d0a1b3dd5202c02d_qw35sft2_f6e61ce4(env, config: dict):
    """Read Sheet2 from SalesRep.xlsx expecting Month, Total, Average columns."""
    file_path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        result = {'header_a1': ws['A1'].value, 'header_b1': ws['B1'].value, 'header_c1': ws['C1'].value, 'months': [], 'totals': [], 'averages': []}
        for row in range(2, 8):
            m = ws.cell(row=row, column=1).value
            t = ws.cell(row=row, column=2).value
            a = ws.cell(row=row, column=3).value
            result['months'].append(str(m).strip() if m else None)
            result['totals'].append(float(t) if t is not None else None)
            result['averages'].append(float(a) if a is not None else None)
        return result
    finally:
        os.unlink(tmp_path)

def get_calc_c1_d1_state__3be2ab82396c7145b9caded82bed3999_qw35sft2_77ca4cb4(env, config: dict):
    """Get C1 number format, D1 formula and D1 computed value from the spreadsheet."""
    path = config.get('path', '/home/user/Padding_Decimals_In_Formular.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb_f = openpyxl.load_workbook(tmp_path, data_only=False)
        ws_f = wb_f.worksheets[0]
        d1_formula = ws_f['D1'].value
        c1_number_format = ws_f['C1'].number_format
        wb_v = openpyxl.load_workbook(tmp_path, data_only=True)
        ws_v = wb_v.worksheets[0]
        d1_value = ws_v['D1'].value
        return {'c1_number_format': c1_number_format, 'd1_formula': str(d1_formula) if d1_formula is not None else '', 'd1_value': str(d1_value) if d1_value is not None else ''}
    finally:
        os.unlink(tmp_path)

def get_ramp_accel_cells__f874da4b0492e2e023290240d31f1c8d_qw35sft2_32041586(env, config: dict):
    """Read key acceleration cell values from RampUpAndDown.xlsx for columns B and D."""
    import tempfile
    import os
    import openpyxl
    file_bytes = env.controller.get_file(config.get('path', '/home/user/RampUpAndDown.xlsx'))
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'b10': ws['B10'].value, 'b30': ws['B30'].value, 'd10': ws['D10'].value, 'd30': ws['D30'].value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_pct_format__6c9c981a9ca6d475ebfb8637fc973b0b_qw35sft2_4e1f3e8b(env, config: dict):
    """Read Sheet2 values and number formats for B2:D6 from SmallBalanceSheet.xlsx."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/SmallBalanceSheet.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        headers = [ws.cell(row=1, column=c).value for c in range(1, 5)]
        data_rows = []
        for r in range(2, 7):
            row_vals = [ws.cell(row=r, column=c).value for c in range(1, 5)]
            data_rows.append(row_vals)
        formats = []
        for r in range(2, 7):
            row_fmts = [ws.cell(row=r, column=c).number_format for c in range(2, 5)]
            formats.append(row_fmts)
        return {'headers': headers, 'data_rows': data_rows, 'formats': formats}
    finally:
        os.unlink(tmp_path)

def get_xlsx_padded_max__7fad2dd93f80fb906237dbbf8bdc9fb2_qw35sft2_5fd89774(env, config: dict):
    """Read D column zero-padded values and E1 cell for max Old ID."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/Customers_New_7digit_Id.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        rows_data = []
        for row in range(2, 31):
            c_val = ws.cell(row=row, column=3).value
            d_val = ws.cell(row=row, column=4).value
            if c_val is not None:
                rows_data.append({'c': int(c_val), 'd': str(d_val) if d_val is not None else None})
        e1_val = ws.cell(row=1, column=5).value
        return {'rows': rows_data, 'e1': e1_val}
    finally:
        os.unlink(tmp_path)

def get_xlsx_spent_and_date_format__2b619f85a4e6da743c7b9581de6419ad_qw35sft2_f6d401cf(env, config: dict):
    """Read C2:C8 number formats and B2:B8 date formats from Keep_Two_decimal_points.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Keep_Two_decimal_points.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        c_formats = [ws.cell(row=r, column=3).number_format for r in range(2, 9)]
        b_formats = [ws.cell(row=r, column=2).number_format for r in range(2, 9)]
        return {'c_formats': c_formats, 'b_formats': b_formats}
    finally:
        os.unlink(tmp_path)

def get_calc_pivot_and_sort__6bb297a00c63757f69bb3bc219190d5b_qw35sft2_0f9cdc1c(env, config: dict):
    """
    Get pivot table existence in Sheet2 and first data row Sales value in Sheet1
    after expected descending sort by Sales.
    Returns: sheet2_exists, sheet1_g2_value (first sale after sort, should be 750)
    """
    file_bytes = env.controller.get_file('/home/user/Invoices.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        sheet2_exists = 'Sheet2' in sheet_names
        invoice_count_10505 = None
        if sheet2_exists:
            ws2 = wb['Sheet2']
            for row in ws2.iter_rows(min_row=1, max_row=ws2.max_row, min_col=1, max_col=2, values_only=True):
                a_val, b_val = row
                try:
                    if int(a_val) == 10505:
                        invoice_count_10505 = b_val
                        break
                except (TypeError, ValueError):
                    pass
        sheet1_g2 = None
        if 'Sheet1' in sheet_names:
            ws1 = wb['Sheet1']
            sheet1_g2 = ws1['G2'].value
        return {'sheet2_exists': sheet2_exists, 'invoice_count_10505': invoice_count_10505, 'sheet1_g2': sheet1_g2}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet1_cell__13e2a9cff490c6125c903160099b9b7a_qw35sft2_86b16597(env, config: dict):
    """Read Sheet1!I1 and Sheet2 pivot table data from SummerSales.xlsx.

    Returns a dict with:
      - value: the string in Sheet1!I1
      - sheet2_exists: bool
      - sheet2_nonempty: count of non-empty cells in Sheet2
      - sheet2_str_values: list of string representations of non-empty Sheet2 cell values
    """
    path = config.get('path', '/home/user/SummerSales.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        i1_value = None
        if 'Sheet1' in wb.sheetnames:
            i1_value = wb['Sheet1']['I1'].value
        sheet2_exists = 'Sheet2' in wb.sheetnames
        sheet2_nonempty = 0
        sheet2_str_values = []
        if sheet2_exists:
            ws2 = wb['Sheet2']
            for row in ws2.iter_rows():
                for cell in row:
                    if cell.value is not None and str(cell.value).strip():
                        sheet2_nonempty += 1
                        sheet2_str_values.append(str(cell.value).strip())
        return {'value': i1_value, 'sheet2_exists': sheet2_exists, 'sheet2_nonempty': sheet2_nonempty, 'sheet2_str_values': sheet2_str_values}
    finally:
        os.unlink(tmp_path)

def get_calc_sort_sum__8b987ab7e49940ce5b75cd3329aaf643_qw35sft2_46f169cb(env, config: dict):
    """Get sorted order indicator and G1 cell value from BoomerangSales.xlsx."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/BoomerangSales.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        dates = []
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=1).value
            if val is not None:
                dates.append(val)
        a2_value = ws.cell(row=2, column=1).value
        is_sorted = len(dates) > 0 and a2_value == min(dates)
        g1_value = ws.cell(row=1, column=7).value
        return {'is_sorted_ascending': is_sorted, 'g1_value': g1_value, 'row_count': ws.max_row - 1}
    finally:
        os.unlink(tmp_path)

def get_xlsx_passfail_and_validation__1231c66fffc58f9761299238f8444bbb_qw35sft2_02b61c7f(env, config: dict):
    """Read D2:D29 values and check data validation from Order_Id_Mark_Pass_Fail.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Order_Id_Mark_Pass_Fail.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        d_values = {}
        for row in range(2, 30):
            d_values[f'D{row}'] = ws.cell(row=row, column=4).value
        has_list_validation = False
        validation_covers_d2_d29 = False
        for dv in ws.data_validations.dataValidation:
            if dv.type == 'list' and dv.formula1:
                sqref_str = str(dv.sqref)
                if 'D2' in sqref_str or 'D:D' in sqref_str:
                    has_list_validation = True
                    formula = dv.formula1.strip('"\'')
                    entries = set((e.strip() for e in formula.replace('\n', ',').split(',') if e.strip()))
                    if {'Pass', 'Fail', 'Held'}.issubset(entries) or entries == {'Pass', 'Fail', 'Held'}:
                        validation_covers_d2_d29 = True
        return {'d_values': d_values, 'has_list_validation': has_list_validation, 'validation_correct': validation_covers_d2_d29}
    finally:
        os.unlink(tmp_path)

def get_month_pivot_sheet2__360bd21f4036eac9068c6b83aba7e2e2_qw35sft2_faf7f287(env, config: dict):
    """Read Sheet2 to verify monthly revenue summary was created."""
    file_path = config.get('path', '/home/user/EntireSummerSales.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        if 'Sheet2' not in sheet_names:
            return {'error': 'Sheet2 not found', 'sheet_names': sheet_names}
        ws2 = wb['Sheet2']
        numeric_values = []
        string_values = []
        for row in ws2.iter_rows(values_only=True):
            for v in row:
                if v is None:
                    continue
                if isinstance(v, (int, float)):
                    numeric_values.append(float(v))
                elif isinstance(v, str):
                    string_values.append(v)
        return {'sheet2_exists': True, 'numeric_values': numeric_values, 'string_values': string_values}
    finally:
        os.unlink(tmp_path)

def get_freeze_and_sheet_name__c058d8a3342bf3e7968eb49c80bec94d_qw35sft2_20f6f927(env, config: dict):
    """Get freeze_panes state and active sheet name from the xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Freeze_row_column.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        freeze_panes = ws.freeze_panes
        sheet_name = ws.title
        return {'freeze_panes': str(freeze_panes) if freeze_panes else None, 'sheet_name': sheet_name}
    finally:
        os.unlink(tmp_path)

def get_sheet1_cell_e1__e7e54b6e523b2a67e0ff77298aeb57dc_qw35sft2_41105dd3(env, config: dict):
    """Read cell E1 value from Sheet1 of DemographicProfile.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/DemographicProfile.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb['Sheet1']
        val = ws['E1'].value
        return {'value': val}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__57a509452fa7aaef388db657e08eeb87_qw35sft2_6006cbc7(env, config: dict):
    """Get a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        return {'value': ws[config.get('cell', 'A1')].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__68e5aa381affaca2d81c1d380acacc93_qw35sft2_b1c6ebcb(env, config: dict):
    """Read a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Create_column_charts_using_statistics.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'value': None}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell_ref = config.get('cell', 'K1')
        value = ws[cell_ref].value
        return {'value': value, 'cell': cell_ref}
    except Exception as e:
        return {'error': str(e), 'value': None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__009b0a2429b483dc78b74a0168744c1a_qw35sft2_3bc353e3(env, config: dict):
    """Read value of a specific cell from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Names_Duplicate_Unique.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell_ref = config.get('cell', 'D2')
        return {'value': ws[cell_ref].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_multi_cells__142de785ea8e1e7c66581a5013a34ff3_qw35sft2_c7bebb74(env, config: dict):
    """Get multiple cell values from a named sheet in an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', '')
        if sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheet_names': wb.sheetnames}
        ws = wb[sheet_name]
        cells = config.get('cells', [])
        values = {}
        for cell_ref in cells:
            values[cell_ref] = ws[cell_ref].value
        return {'values': values, 'sheet_name': sheet_name}
    finally:
        os.unlink(tmp_path)

def get_sheet_names__051a1a86398cf9c50b041fec4af58dd3_qw35sft2_c30ab2c7(env, config: dict):
    """Get the ordered list of sheet names from copy_sheet_insert.xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/copy_sheet_insert.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__b0dcd11a1bdcc1e13016c1ab56479e5f_qw35sft2_43b2aa29(env, config: dict):
    """Read a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell = config.get('cell', 'F10')
        return {'value': ws[cell].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_c7__76b216d4ab5877594f0cff145afca3ba_qw35sft2_fc72e06d(env, config: dict):
    """Read cell C7 value from Represent_in_millions_billions.xlsx."""
    import tempfile
    import os
    import openpyxl
    file_path = config.get('path', '/home/user/Represent_in_millions_billions.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        value = ws['C7'].value
        return {'value': str(value) if value is not None else ''}
    finally:
        os.unlink(tmp_path)

def get_calc_summary_sheet__878e408cd8149dcd33226991d9a64c87_qw35sft2_fa464251(env, config: dict):
    """Read sheet named 'Summary': headers in row 1, sums in row 2."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/NetIncome.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        if 'Summary' not in sheet_names:
            return {'error': 'Summary sheet not found', 'sheets': sheet_names, 'sheet_exists': False}
        ws = wb['Summary']
        return {'sheet_exists': True, 'A1': ws['A1'].value, 'B1': ws['B1'].value, 'A2': ws['A2'].value, 'B2': ws['B2'].value}
    finally:
        os.unlink(tmp_path)

def get_sheet1_net_income_sheet2_revenue__5d938927d392c1d220c2599bb703f167_qw35sft2_9ed717b7(env, config: dict):
    """Get Sheet2 Revenue column state and Sheet1 column C (Net Income) state."""
    file_bytes = env.controller.get_file('/home/user/NetIncome.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        result = {}
        if 'Sheet2' in wb.sheetnames:
            ws2 = wb['Sheet2']
            result['sheet2_header'] = ws2['A1'].value
            result['sheet2_row_count'] = sum((1 for r in range(2, 21) if ws2.cell(row=r, column=1).value is not None))
        else:
            result['sheet2_header'] = None
            result['sheet2_row_count'] = 0
        ws1 = wb['Sheet1']
        result['sheet1_c_header'] = ws1['C1'].value
        c_values = [ws1.cell(row=r, column=3).value for r in range(2, 21)]
        result['sheet1_c_row_count'] = sum((1 for v in c_values if v is not None))
        result['sheet1_c_sum'] = sum((v for v in c_values if isinstance(v, (int, float))))
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_level_primary_range__74bebda311b0ec1306e448288c71ed13_qw35sft2_94755150(env, config: dict):
    """Read B3:B6 cell values from Student_Level_Fill_Blank.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/Student_Level_Fill_Blank.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'B3': ws['B3'].value, 'B4': ws['B4'].value, 'B5': ws['B5'].value, 'B6': ws['B6'].value}
    finally:
        os.unlink(tmp_path)

def get_calc_cell_b12__dc5d11c569bf87971ba13d9e911e3365_qw35sft2_1bf8d790(env, config: dict):
    """Read cell B12 value from WeeklySales.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/WeeklySales.xlsx')
    if not file_bytes:
        return {'error': 'File not found', 'value': None}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        return {'value': ws['B12'].value}
    finally:
        os.unlink(tmp_path)

def get_calc_multi_cells__2412d7535e20b0348d2c25accb946746_qw35sft2_de3f29eb(env, config: dict):
    """Read values of cells E3 and E4 from the saved xlsx file."""
    path = config.get('path', '/home/user/Multiply_Time_Number.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        e3 = ws['E3'].value
        e4 = ws['E4'].value
        return {'e3_value': float(e3) if e3 is not None else None, 'e4_value': float(e4) if e4 is not None else None}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_colors__1a5d60ee39c11d25031d40117b38947f_qw35sft2_ceaed29a(env, config: dict):
    """Get background fill colors for specified cells in an xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Calendar_Highlight_Weekend_Days.xlsx')
    cells_to_check = config.get('cells', [])
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for cell_ref in cells_to_check:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and (fill.fgColor.type == 'rgb'):
                color = fill.fgColor.rgb
            else:
                color = '00000000'
            result[cell_ref] = color
        return result
    finally:
        os.unlink(tmp_path)

def get_sheet_names__0c20409273ae865247912386f480b246_qw35sft2_97f98e80(env, config: dict):
    """Get sheet names from the xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Set_Decimal_Separator_Dot.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_transposed_sorted__f581e5220b9d6fb8293dc4e0f669a755_qw35sft2_eb9648e1(env, config: dict):
    """Read student names and marks from transposed+sorted rows B9:E12 in Students_Class_Subject_Marks.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Students_Class_Subject_Marks.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'B8': ws['B8'].value, 'B9': ws['B9'].value, 'E9': ws['E9'].value, 'B10': ws['B10'].value, 'E10': ws['E10'].value, 'B11': ws['B11'].value, 'E11': ws['E11'].value, 'B12': ws['B12'].value, 'E12': ws['E12'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_zone1_row_totals__3aa70d99f23178c4e0eb339ad1e35c2e_qw35sft2_5d914765(env, config: dict):
    """Read Zone 1 product row totals (F3, F4, F5) from Quarterly_Product_Sales_by_Zone.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Quarterly_Product_Sales_by_Zone.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'F3': ws['F3'].value, 'F4': ws['F4'].value, 'F5': ws['F5'].value}
    finally:
        os.unlink(tmp_path)

def get_sheet2_with_grand_total__0378b78b340823f75ef2df5ebaa9b121_qw35sft2_2dff0f04(env, config: dict):
    """Read Sheet2 from SalesRep.xlsx and extract headers, monthly data, and grand total row."""
    file_path = config.get('path', '/home/user/SalesRep.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        result = {'header_a1': ws['A1'].value, 'header_b1': ws['B1'].value, 'months': [], 'totals': [], 'row8_label': ws['A8'].value, 'row8_value': ws['B8'].value}
        for row in range(2, 8):
            result['months'].append(ws.cell(row=row, column=1).value)
            raw = ws.cell(row=row, column=2).value
            result['totals'].append(float(raw) if raw is not None else None)
        return result
    finally:
        os.unlink(tmp_path)

def get_feb_max_cell__663a803da15ba1a41afdeb0114613f5e_qw35sft2_41dfa63d(env, config: dict):
    """Read cell D24 from the OrderId_Month_Chart.xlsx file to get the Feb max formula result."""
    file_path = config.get('path', '/home/user/OrderId_Month_Chart.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cell_value = ws['D24'].value
        return {'value': cell_value}
    finally:
        os.unlink(tmp_path)

def get_calc_bold_header__771ce45281d75f3370a131f92752ad7a_qw35sft2_d69ed0dd(env, config: dict):
    """Read C2, C3 values and A1, C1 bold formatting from the movie titles spreadsheet."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Movie_title_garbage_clean.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        c2_value = ws['C2'].value
        c3_value = ws['C3'].value
        c1_font = ws['C1'].font
        c1_bold = bool(c1_font.bold) if c1_font else False
        a1_font = ws['A1'].font
        a1_bold = bool(a1_font.bold) if a1_font else False
        return {'c2_value': c2_value, 'c3_value': c3_value, 'c1_bold': c1_bold, 'a1_bold': a1_bold}
    finally:
        os.unlink(tmp_path)

def get_calc_d1_value__62d484d2d5eca5e0b713e4dbdcc095e8_qw35sft2_08e8e0f0(env, config: dict):
    """Get D1 computed value and formula from the spreadsheet."""
    path = config.get('path', '/home/user/Padding_Decimals_In_Formular.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb_f = openpyxl.load_workbook(tmp_path, data_only=False)
        ws_f = wb_f.worksheets[0]
        d1_formula = ws_f['D1'].value
        wb_v = openpyxl.load_workbook(tmp_path, data_only=True)
        ws_v = wb_v.worksheets[0]
        d1_value = ws_v['D1'].value
        return {'d1_formula': str(d1_formula) if d1_formula is not None else '', 'd1_value': str(d1_value) if d1_value is not None else ''}
    finally:
        os.unlink(tmp_path)

def get_sheet2_sorted__2fc0837f7c43e49c220e69070cca644b_qw35sft2_b844fb23(env, config: dict):
    """Read Sheet2 year order and values from SmallBalanceSheet.xlsx."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/SmallBalanceSheet.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        headers = [ws.cell(row=1, column=c).value for c in range(1, 5)]
        data_rows = []
        for r in range(2, 7):
            row_vals = [ws.cell(row=r, column=c).value for c in range(1, 5)]
            data_rows.append(row_vals)
        return {'headers': headers, 'data_rows': data_rows}
    finally:
        os.unlink(tmp_path)

def get_xlsx_padded_count__50235d7deccc7f57d815dc66d0e983bc_qw35sft2_d9297992(env, config: dict):
    """Read D column zero-padded values and E1 customer count from workbook."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/Customers_New_7digit_Id.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        rows_data = []
        for row in range(2, 31):
            c_val = ws.cell(row=row, column=3).value
            d_val = ws.cell(row=row, column=4).value
            if c_val is not None:
                rows_data.append({'c': int(c_val), 'd': str(d_val) if d_val is not None else None})
        e1_val = ws.cell(row=1, column=5).value
        return {'rows': rows_data, 'e1': e1_val}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_bold__636602487d408f3cbef02e77b45f21ff_qw35sft2_c0c36e3c(env, config: dict):
    """Read bold formatting for header cells C4:H4 and sheet zoom scale
    from Zoom_Out_Oversized_Cells.xlsx on VM.

    Returns a dict with:
      - 'C4'..'H4': True if cell font.bold is explicitly True, else None/False
      - 'zoom_scale': integer zoom percentage stored in the sheet view
                      (0 means LibreOffice default, i.e. 100%)
    """
    file_bytes = env.controller.get_file('/home/user/Zoom_Out_Oversized_Cells.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        cells = ['C4', 'D4', 'E4', 'F4', 'G4', 'H4']
        result = {cell: ws[cell].font.bold for cell in cells}
        zoom = ws.sheet_view.zoomScale if ws.sheet_view else 0
        result['zoom_scale'] = zoom if zoom else 0
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_sort_and_sum__7b14f7b3fd242abdf7c663f7ccb97562_qw35sft2_d057b900(env, config: dict):
    """Read D2 (first amount after sort) and D20 (expected SUM total) from spreadsheet."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Arrang_Value_min_to_max.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'d2': ws.cell(row=2, column=4).value, 'd20': ws.cell(row=20, column=4).value}
    finally:
        os.unlink(tmp_path)

def get_calc_sort_sumif__bda50c4eee48b4e3a15799921e84da62_qw35sft2_2e4e7995(env, config: dict):
    """Get sorted order indicator and G1 cell value (SUMIF amazon.com quantity) from BoomerangSales.xlsx."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/BoomerangSales.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        dates = []
        for row in range(2, ws.max_row + 1):
            val = ws.cell(row=row, column=1).value
            if val is not None:
                dates.append(val)
        a2_value = ws.cell(row=2, column=1).value
        is_sorted = len(dates) > 0 and a2_value == min(dates)
        g1_value = ws.cell(row=1, column=7).value
        return {'is_sorted_ascending': is_sorted, 'g1_value': g1_value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_spent_format_total_row__40a8f46c1e43b9d32b31b2230ffda8db_qw35sft2_fb219596(env, config: dict):
    """Read C2:C8 formats, A9 label and C9 value from Keep_Two_decimal_points.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Keep_Two_decimal_points.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        c_formats = [ws.cell(row=r, column=3).number_format for r in range(2, 9)]
        a9_value = ws['A9'].value
        c9_value = ws['C9'].value
        return {'c_formats': c_formats, 'a9_value': a9_value, 'c9_value': c9_value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sales_sum_cell__c61f5322a84515967a3c4cc16d8ae654_qw35sft2_13b519fd(env, config: dict):
    """Get the value in cell F24 (expected SUM row after last data row) of the spreadsheet."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Name_Order_Id_move_column.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        sum_value = ws.cell(row=24, column=6).value
        sales_header = ws.cell(row=1, column=6).value
        return {'sum_value': sum_value, 'sales_header': sales_header}
    finally:
        os.unlink(tmp_path)

def get_sheet_name_and_pdf__57d2609933a3290edf4943a154f0921f_qw35sft2_ce5aa5ea(env, config: dict):
    """Check sheet name, page fit-to-one-page scaling, and PDF export of saved xlsx."""
    pdf_bytes = env.controller.get_file('/home/user/Resize_Cells_Fit_Page.pdf')
    pdf_exists = bool(pdf_bytes and len(pdf_bytes) > 0)
    xlsx_path = config.get('path', '/home/user/Resize_Cells_Fit_Page.xlsx')
    file_bytes = env.controller.get_file(xlsx_path)
    if not file_bytes:
        return {'error': 'xlsx not found', 'pdf_exists': pdf_exists, 'sheet_name': None, 'page_fit': False}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = wb.sheetnames[0] if wb.sheetnames else None
        ws = wb.worksheets[0]
        try:
            page_setup_pr = ws.sheet_properties.pageSetUpPr if ws.sheet_properties else None
            page_fit = bool(page_setup_pr and page_setup_pr.fitToPage)
        except Exception:
            page_fit = False
        return {'pdf_exists': pdf_exists, 'sheet_name': sheet_name, 'page_fit': page_fit}
    finally:
        os.unlink(tmp_path)

def get_calc_pivot_and_total__2c779926d5355737435372c9e97f0d18_qw35sft2_43ca166f(env, config: dict):
    """
    Get pivot table existence in Sheet2 and SUM formula result in Sheet1 cell G20.
    Returns: sheet2_exists, invoice_count_10505 (spot check), sheet1_g20_value
    """
    file_bytes = env.controller.get_file('/home/user/Invoices.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        sheet2_exists = 'Sheet2' in sheet_names
        invoice_count_10505 = None
        if sheet2_exists:
            ws2 = wb['Sheet2']
            for row in ws2.iter_rows(min_row=1, max_row=ws2.max_row, min_col=1, max_col=2, values_only=True):
                a_val, b_val = row
                try:
                    if int(a_val) == 10505:
                        invoice_count_10505 = b_val
                        break
                except (TypeError, ValueError):
                    pass
        sheet1_g20 = None
        if 'Sheet1' in sheet_names:
            ws1 = wb['Sheet1']
            sheet1_g20 = ws1['G20'].value
        return {'sheet2_exists': sheet2_exists, 'invoice_count_10505': invoice_count_10505, 'sheet1_g20': sheet1_g20}
    finally:
        os.unlink(tmp_path)

def get_xlsx_passfail_and_count__9a44e82214177e638fd33343975b4139_qw35sft2_82e0b7b8(env, config: dict):
    """Read D2:D29, E1 (COUNTIF formula result), and data validation from Order_Id_Mark_Pass_Fail.xlsx."""
    file_bytes = env.controller.get_file('/home/user/Order_Id_Mark_Pass_Fail.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        d_values = {}
        for row in range(2, 30):
            d_values[f'D{row}'] = ws.cell(row=row, column=4).value
        e1_value = ws['E1'].value
        e2_value = ws['E2'].value
        has_list_validation = False
        for dv in ws.data_validations.dataValidation:
            if dv.type == 'list' and dv.formula1:
                sqref_str = str(dv.sqref)
                if 'D2' in sqref_str or 'D:D' in sqref_str:
                    formula = dv.formula1.strip('"\'')
                    entries = set((e.strip() for e in formula.replace('\n', ',').split(',') if e.strip()))
                    if {'Pass', 'Fail', 'Held'}.issubset(entries):
                        has_list_validation = True
        return {'d_values': d_values, 'e1_value': e1_value, 'e2_value': e2_value, 'has_list_validation': has_list_validation}
    finally:
        os.unlink(tmp_path)

def get_pivot_and_sorted__3da7df8645c1252be37e4e927311c6ef_qw35sft2_ee96d6d5(env, config: dict):
    """Read Sheet2 numeric values and Sheet1 first 3 data rows Revenue column to check sort order."""
    file_path = config.get('path', '/home/user/EntireSummerSales.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        ws1 = wb['Sheet1'] if 'Sheet1' in sheet_names else wb.worksheets[0]
        sheet1_revenues = []
        for row_idx in range(2, 7):
            v = ws1.cell(row=row_idx, column=7).value
            if v is not None:
                try:
                    sheet1_revenues.append(float(v))
                except (TypeError, ValueError):
                    pass
        sheet2_numeric = []
        if 'Sheet2' in sheet_names:
            ws2 = wb['Sheet2']
            for row in ws2.iter_rows(values_only=True):
                for v in row:
                    if isinstance(v, (int, float)):
                        sheet2_numeric.append(float(v))
        return {'sheet1_top_revenues': sheet1_revenues, 'sheet2_exists': 'Sheet2' in sheet_names, 'sheet2_numeric': sheet2_numeric}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet1_cell__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_dfdd8fd9(env, config: dict):
    """Read cell I1 from Sheet1 of SummerSales.xlsx on VM."""
    path = config.get('path', '/home/user/SummerSales.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet1' not in wb.sheetnames:
            return {'error': 'Sheet1 not found'}
        ws = wb['Sheet1']
        return {'value': ws['I1'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_combined__beec5e7b513d1d21dbf198b76143a06e_qw35sft2_52d26355(env, config: dict):
    """Read Sheet1!I1 and check Sheet2 for product and sales-channel pivot tables."""
    path = config.get('path', '/home/user/SummerSales.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet1_i1 = None
        if 'Sheet1' in wb.sheetnames:
            sheet1_i1 = wb['Sheet1']['I1'].value
        if 'Sheet2' not in wb.sheetnames:
            return {'sheet1_i1': sheet1_i1, 'sheet2_exists': False, 'sheet2_has_product': False, 'sheet2_has_channel': False, 'sheet2_cell_count': 0}
        ws2 = wb['Sheet2']
        sheet2_texts = []
        for row in ws2.iter_rows():
            for cell in row:
                if cell.value is not None:
                    sheet2_texts.append(str(cell.value).strip())
        text_lower = ' '.join(sheet2_texts).lower()
        has_product = 'product' in text_lower
        has_channel = 'channel' in text_lower
        return {'sheet1_i1': sheet1_i1, 'sheet2_exists': True, 'sheet2_has_product': has_product, 'sheet2_has_channel': has_channel, 'sheet2_cell_count': len(sheet2_texts)}
    finally:
        os.unlink(tmp_path)

def get_sheet_names__0ed76c34ce42b49c8d0f69a316afbc19_qw35sft2_1c79a8a7(env, config: dict):
    """Get sheet names from DemographicProfile.xlsx on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/DemographicProfile.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': wb.sheetnames}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_cell__fd1c3b4eab4669bd9848fe12e4f0d1d4_qw35sft2_48f35820(env, config: dict):
    """Get cell value from a named sheet in an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_name = config.get('sheet_name', '')
        if sheet_name not in wb.sheetnames:
            return {'error': f'Sheet "{sheet_name}" not found', 'sheet_names': wb.sheetnames}
        ws = wb[sheet_name]
        cell_ref = config.get('cell', 'A1')
        value = ws[cell_ref].value
        return {'value': value, 'sheet_name': sheet_name, 'cell': cell_ref}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cells_range__1d24482d98b418ee12818e4d70230d5c_qw35sft2_74c8c004(env, config: dict):
    """Read values from a range of cells in an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Names_Duplicate_Unique.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        start_row = config.get('start_row', 2)
        end_row = config.get('end_row', 7)
        col = config.get('col', 4)
        values = []
        for row in range(start_row, end_row + 1):
            values.append(ws.cell(row=row, column=col).value)
        return {'values': values}
    finally:
        os.unlink(tmp_path)

def get_xlsx_student_c3_c4__87a420fb18a4347285a9615e2d7a9d87_qw35sft2_70b598e6(env, config: dict):
    """Read C3 and C4 cell values from Student_Level_Fill_Blank.xlsx on VM."""
    file_bytes = env.controller.get_file('/home/user/Student_Level_Fill_Blank.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'C3': ws['C3'].value, 'C4': ws['C4'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sheet_names__08c531e403541f6bab4f59a478d0e6c2_qw35sft2_01a31a75(env, config: dict):
    """Get the list of sheet names from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        return {'sheet_names': list(wb.sheetnames)}
    finally:
        os.unlink(tmp_path)

def get_sheet_and_cell__ebd5692a1e5abf6c2058f3bfb502b3de_qw35sft2_eb586ff4(env, config: dict):
    """Get sheet names and cell A1 value from the 'LARS Resources (Backup)' sheet."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file('/home/user/copy_sheet_insert.xlsx')
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_names = wb.sheetnames
        backup_a1 = None
        if 'LARS Resources (Backup)' in sheet_names:
            ws = wb['LARS Resources (Backup)']
            backup_a1 = ws['A1'].value
        return {'sheet_names': sheet_names, 'backup_a1': backup_a1}
    finally:
        os.unlink(tmp_path)

def get_calc_sheet2_3col__581360a1329c03f025e39497ca2b0766_qw35sft2_c19c273c(env, config: dict):
    """Read Sheet2 headers and sum values for 3 columns: Total Revenue, Total Expenses, Net Income."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/NetIncome.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'Sheet2' not in wb.sheetnames:
            return {'error': 'Sheet2 not found', 'sheets': wb.sheetnames}
        ws = wb['Sheet2']
        return {'A1': ws['A1'].value, 'B1': ws['B1'].value, 'C1': ws['C1'].value, 'A2': ws['A2'].value, 'B2': ws['B2'].value, 'C2': ws['C2'].value}
    finally:
        os.unlink(tmp_path)

def get_calc_multi_cells__ddba4585b2adadf3e6cdf7efce56b822_qw35sft2_bdd7f58f(env, config: dict):
    """Read values of cells E3 and G3 from the saved xlsx file."""
    path = config.get('path', '/home/user/Multiply_Time_Number.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        e3 = ws['E3'].value
        g3 = ws['G3'].value
        return {'e3_value': float(e3) if e3 is not None else None, 'g3_value': float(g3) if g3 is not None else None}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_colors__cc3c6531bd3eb1e9b2189f82dea73399_qw35sft2_64440f9c(env, config: dict):
    """Get background fill colors for specified cells in an xlsx file."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Calendar_Highlight_Weekend_Days.xlsx')
    cells_to_check = config.get('cells', [])
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        result = {}
        for cell_ref in cells_to_check:
            cell = ws[cell_ref]
            fill = cell.fill
            if fill and fill.fgColor and (fill.fgColor.type == 'rgb'):
                color = fill.fgColor.rgb
            else:
                color = '00000000'
            result[cell_ref] = color
        return result
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__9d8210ef208924a20f5946d708d3b9b8_qw35sft2_fb0b689c(env, config: dict):
    """Read a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Create_column_charts_using_statistics.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'value': None}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        sheet_idx = config.get('sheet', 0)
        ws = wb.worksheets[sheet_idx]
        cell_ref = config.get('cell', 'K3')
        value = ws[cell_ref].value
        return {'value': value, 'cell': cell_ref}
    except Exception as e:
        return {'error': str(e), 'value': None}
    finally:
        os.unlink(tmp_path)

def get_pptx_table_cells__d704f09ae59c96034fafb3c1ef256369_qw35sft2_62936050(env, config: dict):
    """Get the first row and a specific cell from a table on slide 4 of a PPTX file."""
    try:
        from pptx import Presentation
        file_bytes = env.controller.get_file(config.get('path', '/home/user/Desktop/33_1.pptx'))
        if not file_bytes:
            return {'error': 'File not found'}
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            prs = Presentation(tmp_path)
            slide = prs.slides[3]
            for shape in slide.shapes:
                if shape.shape_type == 19:
                    table = shape.table
                    row0 = [cell.text.strip() for cell in table.rows[0].cells]
                    cell_row = config.get('cell_row', 1)
                    cell_col = config.get('cell_col', 0)
                    specific_cell = table.rows[cell_row].cells[cell_col].text.strip()
                    return {'table_row0': row0, 'specific_cell': specific_cell}
            return {'error': 'No table found on slide 4'}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_docx_table_last_cell__818c616cbd97d0d5bf2ff81918416501_qw35sft2_fcef84bd(env, config: dict):
    """Download Table_Of_Work_Effort_Instructions.docx, return table count, dims, and last cell text."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_rows': 0, 'last_cols': 0, 'last_cell': ''}
        last = tables[-1]
        rows = len(last.rows)
        cols = len(last.columns)
        last_cell = last.rows[-1].cells[-1].text.strip() if rows > 0 and cols > 0 else ''
        return {'table_count': table_count, 'last_rows': rows, 'last_cols': cols, 'last_cell': last_cell}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_docx_table_with_cell__2aa481b876a15405aedb5c91a8fc78e9_qw35sft2_0860cc06(env, config: dict):
    """Download Table_Of_Work_Effort_Instructions.docx, return last table dims and first cell text."""
    import tempfile
    import os
    from docx import Document
    path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'table_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables = doc.tables
        table_count = len(tables)
        if table_count == 0:
            return {'table_count': 0, 'last_rows': 0, 'last_cols': 0, 'cell_a1': ''}
        last = tables[-1]
        rows = len(last.rows)
        cols = len(last.columns)
        cell_a1 = last.rows[0].cells[0].text.strip() if rows > 0 and cols > 0 else ''
        return {'table_count': table_count, 'last_rows': rows, 'last_cols': cols, 'cell_a1': cell_a1}
    except Exception as e:
        return {'error': str(e)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_val__2feaa23241a59b8898b31056b9df1f85_qw35sft2_19216087(env, config: dict):
    """Read a specific cell value from the Employee Performance Evaluation Summary xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        cell_ref = config.get('cell', 'A1')
        value = ws[cell_ref].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_freeze_panes__b12f1d02d5521f031b94e1747a1c556a_qw35sft2_1a028055(env, config: dict):
    """Read freeze_panes setting from the xlsx file at config['path']."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path)
        ws = wb.worksheets[0]
        freeze = ws.freeze_panes
        return {'freeze_panes': str(freeze) if freeze else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_year_column__c99ee3ea1307143970f7d9489f51bcb4_qw35sft2_99620a62(env, config: dict):
    """Download the spreadsheet and read the Year column (column E, rows 1-6)."""
    import tempfile
    import os
    try:
        import openpyxl
    except ImportError:
        return {'error': 'openpyxl not available', 'header': None, 'years': []}
    path = config.get('path', '/home/user/Desktop/rsc-ebook-collection-2023.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found', 'header': None, 'years': []}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        header = ws.cell(row=1, column=5).value
        years = []
        for row_idx in range(2, 7):
            val = ws.cell(row=row_idx, column=5).value
            if val is not None:
                try:
                    years.append(int(val))
                except (ValueError, TypeError):
                    years.append(str(val))
            else:
                years.append(None)
        return {'header': header, 'years': years}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__5777a0629a4e670ae915b1fe64e58378_qw35sft2_b06b1277(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        return {'value': ws[config.get('cell', 'A1')].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_2017_2018_cities__1272197382489180193ce8b35860bb11_qw35sft2_66402127(env, config: dict):
    """Read 2017 and 2018 conference city cells (C14-C19) from ConferenceCity.xlsx."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/Desktop/ConferenceCity.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'c14': ws['C14'].value, 'c15': ws['C15'].value, 'c16': ws['C16'].value, 'c17': ws['C17'].value, 'c18': ws['C18'].value, 'c19': ws['C19'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_val__61d78e261ca58a364d140d5219ceb5c7_qw35sft2_5a873ee6(env, config: dict):
    """Read a specific cell value from the Employee Performance Evaluation Summary xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        cell_ref = config.get('cell', 'A1')
        value = ws[cell_ref].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_header_align__7d3ab32fb77b1ea64cfabfdcd56b69aa_qw35sft2_ea7fbc3e(env, config: dict):
    """Read horizontal alignment of cell A1 from the xlsx file at config['path']."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path)
        ws = wb.worksheets[0]
        alignments = []
        for col in range(1, 5):
            cell = ws.cell(row=1, column=col)
            align = cell.alignment.horizontal if cell.alignment else None
            alignments.append(align)
        return {'header_alignments': alignments, 'a1_alignment': alignments[0] if alignments else None}
    finally:
        os.unlink(tmp_path)

def get_xlsx_sampled_conf_cities__fb22e54de099b00f4d6a127b12703f11_qw35sft2_e8fe6de6(env, config: dict):
    """Read 3 sampled city cells (2013 ICML, 2016 NeurIPS, 2019 ICLR) from ConferenceCity.xlsx."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/Desktop/ConferenceCity.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'c3': ws['C3'].value, 'c13': ws['C13'].value, 'c20': ws['C20'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_row_data__5abc6535b144407839fbbe3d8a497678_qw35sft2_983a349c(env, config: dict):
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        row = config.get('row', 6)
        result = {}
        for col in range(1, ws.max_column + 1):
            header = ws.cell(row=1, column=col).value
            value = ws.cell(row=row, column=col).value
            if header:
                result[header] = value
        return result
    finally:
        os.unlink(tmp_path)

def get_docx_and_xlsx__25c4c78ad235f6f7dec2a3ce1f7f2c6f_qw35sft2_6c93d383(env, config: dict):
    """Read Answer.docx text and a specific xlsx cell value."""
    import tempfile
    import os
    full_text = ''
    docx_error = None
    try:
        from docx import Document
        docx_path = config.get('path', '/home/user/Desktop/Answer.docx')
        file_bytes = env.controller.get_file(docx_path)
        if file_bytes:
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                doc = Document(tmp_path)
                full_text = '\n'.join((p.text for p in doc.paragraphs))
            except Exception as e:
                docx_error = str(e)
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass
        else:
            docx_error = 'Answer.docx not found'
    except ImportError:
        docx_error = 'python-docx not installed'
    cell_value = None
    xlsx_error = None
    try:
        import openpyxl
        xlsx_path = config.get('xlsx_path', '/home/user/Desktop/Course Timetable.xlsx')
        xlsx_cell = config.get('xlsx_cell', 'B9')
        xlsx_bytes = env.controller.get_file(xlsx_path)
        if xlsx_bytes:
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp2:
                tmp2.write(xlsx_bytes)
                tmp2_path = tmp2.name
            try:
                wb = openpyxl.load_workbook(tmp2_path, data_only=True)
                ws = wb.worksheets[0]
                cell_value = ws[xlsx_cell].value
            except Exception as e:
                xlsx_error = str(e)
            finally:
                try:
                    os.unlink(tmp2_path)
                except Exception:
                    pass
        else:
            xlsx_error = 'xlsx not found'
    except ImportError:
        xlsx_error = 'openpyxl not installed'
    return {'full_text': full_text, 'cell_value': cell_value, 'docx_error': docx_error, 'xlsx_error': xlsx_error}

def get_xlsx_cell_val__9d1e0f5a79d29ecfac0c206e85fa98e9_qw35sft2_9aabb6ba(env, config: dict):
    """Read a specific cell value from the Employee Performance Evaluation Summary xlsx."""
    import tempfile, os, openpyxl
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        cell_ref = config.get('cell', 'A1')
        value = ws[cell_ref].value
        return {'value': value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_unseen_sheet_headers__ac461f3f5169eddd647d2f7d2aad85a2_qw35sft2_07dfbad5(env, config: dict):
    """Get the header row of the 'unseen_movies' sheet from the movies.xlsx file."""
    import tempfile
    import os
    import openpyxl
    path = config.get('path', '/home/user/Desktop/movies.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if 'unseen_movies' not in wb.sheetnames:
            return {'error': 'Sheet unseen_movies not found', 'sheets': wb.sheetnames}
        ws = wb['unseen_movies']
        headers = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
        return {'headers': headers, 'sheet_count': len(wb.sheetnames)}
    finally:
        os.unlink(tmp_path)

def get_xlsx_icml_cities__e73f742a9bd3e72f68e6861886fdd0c7_qw35sft2_1ea5a37e(env, config: dict):
    """Read all 7 ICML city cells (rows 3,6,9,12,15,18,21) from ConferenceCity.xlsx."""
    import tempfile, os, openpyxl
    file_path = config.get('path', '/home/user/Desktop/ConferenceCity.xlsx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[0]
        return {'c3': ws['C3'].value, 'c6': ws['C6'].value, 'c9': ws['C9'].value, 'c12': ws['C12'].value, 'c15': ws['C15'].value, 'c18': ws['C18'].value, 'c21': ws['C21'].value}
    finally:
        os.unlink(tmp_path)

def get_xlsx_cell_value__9d415d4db27ab3e1c21d0be924167f0e_qw35sft2_aff163a8(env, config: dict):
    """Read a single cell value from an xlsx file on the VM."""
    import tempfile, os, openpyxl
    path = config.get('path', '/home/user/Desktop/Course Timetable.xlsx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'error': 'File not found'}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.worksheets[config.get('sheet', 0)]
        cell_addr = config.get('cell', 'D3')
        value = ws[cell_addr].value
        return {'value': str(value) if value is not None else None}
    finally:
        os.unlink(tmp_path)
