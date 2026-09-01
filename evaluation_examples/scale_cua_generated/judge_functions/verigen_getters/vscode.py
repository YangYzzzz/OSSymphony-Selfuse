"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_vscode_terminal_status__64ba4345', 'get_vscode_git_exclude__3826e4daac10a542e57e7cc4a2b6b258', 'get_vscode_settings__7a482b52', 'get_vscode_file_count__7c2e5da9', 'get_vscode_open_file__149df8b642d00f3a86d140b4582f20cc', 'get_vscode_file_content__09b11b45', 'get_vscode_settings__ab83ccf1', 'get_vscode_settings__aec709e9', 'get_vscode_settings_json__51833a76bce91777752996df0b22c33d', 'get_vscode_text_content__e99e7d69ea502440f7bd1b2cb57a309d', 'get_multiple_workspace_files__f54bf781c32dea93b80d6edc50522687', 'get_vscode_open_file__c60120ed87e101802d8b786344162dfa', 'get_vscode_adjacent_lines_match__0737d111', 'get_vscode_file_content__57096d124fd446ca6cba7d1316aac1bd', 'get_docx_unique_workspace_count__7f14f5d18ace26bd462391619642ff63', 'get_vscode_file_content__59124b969e280d0086072fa0af69d61b', 'get_vscode_settings__b07969a1', 'get_vscode_autosave__c5efe7eb60d4b82c4d217326823a4934', 'get_vscode_settings__4269362e9e26c329da428a21ddabb7ea', 'get_vscode_tabsize__a6b4939a2d71540bebc75ede5d00109a', 'get_vscode_markdown_files__77849053', 'get_vscode_workspace_files__4ab817c3', 'get_vscode_file_content__b5e3afd36096fa6a3d39a615d4795ccf', 'get_vscode_nodemod_exclude__bdfa2ae4136b63f287ad1f185959812f', 'get_vscode_wordwrap__ca4beaed446890c3b888d4a40b939831', 'get_vscode_line_is_comment__d6c7b90d', 'get_vscode_open_file__005f6ac9', 'get_vscode_settings__a46d28041c0fe91cc8b08a2032679781', 'get_vscode_settings__f8298be8', 'get_vscode_line_contains__b5d41ccd', 'get_vscode_file_content__da119c1e', 'get_vscode_workspace_json__6ed0a554', 'get_vscode_settings__44a94327', 'get_vscode_text_content__66cc1bc44690547654db9625c97259f8', 'get_vscode_workspace__6efce4eb1d6563b7b443059d330de6aa', 'get_vscode_file_content__a5163ad1', 'get_vscode_settings__232783f17230c635388a2d96b7097023', 'get_vscode_config_files__17da0621', 'get_vscode_open_tabs__32a2204c', 'get_vscode_text_content__71cb5a9efacfe89bccff2081bcadcf02', 'get_vscode_text_content__56d0ef227fc0a081b24201d3ecb3d358', 'get_vscode_settings__ff65f872', 'get_vscode_function_docstring__c84ddf7a', 'get_docx_unique_workspace_records__1aff24635b2ad5253401338faad2d9bc', 'get_vscode_file_content__baf14747', 'get_vscode_text_content__681a97c35eb849e5cf9422adfe4e1aea', 'get_vscode_text_content__65dee18f1880f7d028c2b1727fd62d90', 'get_vscode_settings__37a81987aef2f239d369b6635778b9ec', 'get_vscode_extension_dir__42aa293803886f1362c8d3d5e33a2703', 'get_vscode_workspace_opened__56d366f94894cbfe50f3e560644c1cc7', 'get_vscode_settings_autosave__06d1442bba585f1fb8ea50a4956df674', 'get_vscode_file_content__06f73a877cbecd1e2a247622c3136810', 'get_vscode_text_content__62f7a00e7dd5e4db70eb00a615f012ef', 'get_vscode_file_content__219e3148d77134036e282c6dd4b41d12', 'get_vscode_file_line__8f9a3936', 'get_vscode_line_count__b2b950a6', 'get_vscode_search_result__0da38d8e', 'get_vscode_file_content__39395653', 'get_vscode_file_content__75a1331b', 'get_vscode_text_content__c5c5b95c23a9d2c28863362320aba24b', 'get_vscode_text_content__0f72ff3b9c8d9680b46915659d675f48', 'get_vscode_variable_usage__f33f7a7e', 'get_vscode_open_file__ddd2693d', 'get_vscode_open_file__e925f44c', 'get_vscode_line_indentation__4d2e0474', 'get_vscode_file_content__02be85da', 'get_vscode_text_content__dd369016ed084e5c5f565139bb4ef07f', 'get_vscode_lines_order__2c3b878a', 'get_vscode_python_path__1f6f3af3ee1e7e72d7b32984543de005', 'get_vscode_settings_content__a74ebb334721e27ebd0b52511b5b5b07', 'get_vscode_settings__d0166485dd9dee18fc251c7dbffffc79', 'get_vscode_settings__052d7bff3ecc5844bf94710602cff931', 'get_vscode_file_content__9d7d44f7', 'get_vscode_explorer__71b2b2791188a1049bb680b7c56ffb2d', 'get_vscode_file_content__2fa9f229', 'get_vscode_settings__4df3f2e6', 'get_vscode_python_path__e7a35ed0ec20ca7ab6b257f3f5c87e23']

def get_vscode_terminal_status__64ba4345(env, config: Dict[str, Any]) -> str:
    """Get the terminal status from VS Code.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file containing terminal status
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_git_exclude__3826e4daac10a542e57e7cc4a2b6b258(env, config: Dict[str, Any]) -> str:
    """Get VS Code settings.json file content for .git folder exclusion check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded settings.json file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_vscode_settings__7a482b52(env, config: Dict[str, Any]) -> str:
    """Get the VS Code settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file containing VS Code settings
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_count__7c2e5da9(env, config: Dict[str, Any]) -> str:
    """Get the count of specific file types in VSCode workspace."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_open_file__149df8b642d00f3a86d140b4582f20cc(env, config: Dict[str, Any]) -> str:
    """Get the currently open file in VSCode using extension command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, and dest

    Returns:
        str: Path to the file containing the open file info
    """
    from ..getters.file import get_vm_file
    from ..getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_content__09b11b45(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_settings__ab83ccf1(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for font size checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_settings__aec709e9(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for auto-save checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_settings_json__51833a76bce91777752996df0b22c33d(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract settings from VS Code settings.json file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        dict: The parsed JSON settings, or empty dict if file not found or invalid
    """
    file_path = config.get('path', '')
    if not file_path:
        return {}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {}
    try:
        with tempfile.NamedTemporaryFile(suffix='.json', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with open(tmp_path, 'r') as f:
                settings = json.load(f)
            return settings if isinstance(settings, dict) else {}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {}

def get_vscode_text_content__e99e7d69ea502440f7bd1b2cb57a309d(env, config):
    """Extract text content from a file for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the text file, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file {file_path}: {e}')
        return ''

def get_multiple_workspace_files__f54bf781c32dea93b80d6edc50522687(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if multiple workspace files exist on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'workspace_paths' key (list of paths)

    Returns:
        Dict with 'workspaces' list containing info about each workspace
    """
    workspace_paths = config.get('workspace_paths', [])
    workspaces = []
    for path in workspace_paths:
        workspace_info = {'path': path, 'exists': False, 'valid_json': False, 'folders': []}
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            workspace_info['exists'] = True
            try:
                workspace_content = file_bytes.decode('utf-8')
                workspace_data = json.loads(workspace_content)
                workspace_info['valid_json'] = True
                workspace_info['folders'] = workspace_data.get('folders', [])
            except:
                workspace_info['valid_json'] = False
        workspaces.append(workspace_info)
    return {'workspaces': workspaces, 'total_count': len(workspaces), 'exists_count': sum((1 for w in workspaces if w['exists'])), 'valid_count': sum((1 for w in workspaces if w['valid_json']))}

def get_vscode_open_file__c60120ed87e101802d8b786344162dfa(env, config: Dict[str, Any]) -> str:
    """Get the currently open file in VSCode using extension command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, and dest

    Returns:
        str: Path to the file containing the open file info
    """
    from ..getters.file import get_vm_file
    from ..getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_adjacent_lines_match__0737d111(env, config: Dict[str, Any]) -> bool:
    """Check if two adjacent lines have the same content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'line_number'

    Returns:
        True if line matches the next line, False otherwise
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    line_number = config.get('line_number', 0)
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        if line_number < len(lines) - 1:
            line1 = lines[line_number].strip()
            line2 = lines[line_number + 1].strip()
            return line1 == line2 and len(line1) > 0
        return False

def get_vscode_file_content__57096d124fd446ca6cba7d1316aac1bd(env, config):
    """Read text file content from VM for indentation validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_docx_unique_workspace_count__7f14f5d18ace26bd462391619642ff63(env, config):
    """Extract the count that the user wrote at the end of the docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (e.g., {'path': '/home/user/file.docx'})

    Returns:
        int: Count that the user wrote at the end of the document, or 0 if not found
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return 0
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        last_text = None
        for para in reversed(doc.paragraphs):
            text = para.text.strip()
            if text:
                last_text = text
                break
        if not last_text:
            return 0
        import re
        numbers = re.findall('\\b\\d+\\b', last_text)
        if numbers:
            return int(numbers[-1])
        return 0
    except Exception as e:
        logger.error(f'Error processing docx file: {e}')
        return 0
    finally:
        os.unlink(tmp_path)

def get_vscode_file_content__59124b969e280d0086072fa0af69d61b(env, config):
    """Read text file content from VM for indentation validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_vscode_settings__b07969a1(env, config: Dict[str, Any]) -> str:
    """Get the VS Code settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file containing VS Code settings
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_autosave__c5efe7eb60d4b82c4d217326823a4934(env, config: Dict[str, Any]) -> str:
    """Get VS Code settings.json file content for auto-save configuration check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded settings.json file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_vscode_settings__4269362e9e26c329da428a21ddabb7ea(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json content as a dictionary.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dictionary containing the parsed settings.json content, or empty dict on error
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {}
        settings_text = file_bytes.decode('utf-8')
        settings_dict = json.loads(settings_text)
        return settings_dict
    except Exception as e:
        return {}

def get_vscode_tabsize__a6b4939a2d71540bebc75ede5d00109a(env, config: Dict[str, Any]) -> str:
    """Get VS Code settings.json file content for tab size configuration check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded settings.json file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_vscode_markdown_files__77849053(env, config: Dict[str, Any]) -> str:
    """Get markdown files in VSCode workspace."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_workspace_files__4ab817c3(env, config: Dict[str, Any]) -> str:
    """Get the list of workspace files in VSCode via the custom extension command."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_content__b5e3afd36096fa6a3d39a615d4795ccf(env, config):
    """Read text file content from VM for indentation validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_vscode_nodemod_exclude__bdfa2ae4136b63f287ad1f185959812f(env, config: Dict[str, Any]) -> str:
    """Get VS Code settings.json file content for node_modules exclusion check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded settings.json file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_vscode_wordwrap__ca4beaed446890c3b888d4a40b939831(env, config: Dict[str, Any]) -> str:
    """Get VS Code settings.json file content for word wrap configuration check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded settings.json file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_vscode_line_is_comment__d6c7b90d(env, config: Dict[str, Any]) -> bool:
    """Check if a specific line is commented out.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'line_number'

    Returns:
        True if line is a comment, False otherwise
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    line_number = config.get('line_number', 0)
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        if line_number < len(lines):
            line = lines[line_number].lstrip()
            return line.startswith('#')
        return False

def get_vscode_open_file__005f6ac9(env, config: Dict[str, Any]) -> str:
    """Get the currently open file in VSCode via the custom extension command."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_settings__a46d28041c0fe91cc8b08a2032679781(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json content as a dictionary.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dictionary containing the parsed settings.json content, or empty dict on error
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {}
        settings_text = file_bytes.decode('utf-8')
        settings_dict = json.loads(settings_text)
        return settings_dict
    except Exception as e:
        return {}

def get_vscode_settings__f8298be8(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for theme checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_line_contains__b5d41ccd(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a specific line contains a substring with proper indentation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path', 'line_number', and 'substring'

    Returns:
        Dict with 'has_content', 'indentation', and 'reference_indentation' keys
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    line_number = config.get('line_number', 0)
    substring = config.get('substring', '')
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        result = {'has_content': False, 'indentation': '', 'reference_indentation': ''}
        if line_number < len(lines):
            line = lines[line_number]
            if substring in line:
                result['has_content'] = True
                result['indentation'] = line[:len(line) - len(line.lstrip())]
            if line_number + 1 < len(lines):
                ref_line = lines[line_number + 1]
                result['reference_indentation'] = ref_line[:len(ref_line) - len(ref_line.lstrip())]
        return result

def get_vscode_file_content__da119c1e(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_workspace_json__6ed0a554(env, config: Dict[str, Any]) -> Optional[Dict]:
    """
    Get VSCode workspace JSON file content as a dictionary.

    This getter fetches the workspace JSON file from the VM and parses it.

    Args:
        env: Environment object with controller to access VM files
        config: Configuration dict with 'path' key for the workspace file path

    Returns:
        Dict containing the parsed JSON workspace configuration, or None if file cannot be read
    """
    path = config.get('path')
    if not path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get workspace file from VM: {path}')
            return None
        data = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded workspace JSON from {path}')
        return data
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading workspace file {path}: {e}')
        return None

def get_vscode_settings__44a94327(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for tab size checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_text_content__66cc1bc44690547654db9625c97259f8(env, config):
    """Get text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.warning('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_vscode_workspace__6efce4eb1d6563b7b443059d330de6aa(env, config: Dict[str, Any]) -> str:
    """Get the current workspace folder in VSCode using extension command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, and dest

    Returns:
        str: Path to the file containing the workspace info
    """
    from ..getters.file import get_vm_file
    from ..getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_content__a5163ad1(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_settings__232783f17230c635388a2d96b7097023(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VSCode settings from the project.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for settings file location

    Returns:
        Dict containing the settings, or empty dict if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {}
    try:
        settings_json = json.loads(file_bytes.decode('utf-8'))
        return settings_json
    except Exception as e:
        logger.error(f'Error parsing settings.json: {e}')
        return {}

def get_vscode_config_files__17da0621(env, config: Dict[str, Any]) -> str:
    """Get config files in VSCode workspace."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_open_tabs__32a2204c(env, config: Dict[str, Any]) -> str:
    """Get the list of open tabs in VS Code.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file containing open tabs info
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_text_content__71cb5a9efacfe89bccff2081bcadcf02(env, config):
    """Get text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.warning('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_vscode_text_content__56d0ef227fc0a081b24201d3ecb3d358(env, config):
    """Get text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.warning('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_vscode_settings__ff65f872(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for word wrap checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_function_docstring__c84ddf7a(env, config: Dict[str, Any]) -> str:
    """Extract the docstring of a function.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'function_name'

    Returns:
        The docstring if found, empty string otherwise
    """
    from .file import get_vm_file
    import os
    import ast
    path = config['path']
    function_name = config.get('function_name', '')
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        content = f.read()
    try:
        tree = ast.parse(content)
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef) and node.name == function_name:
                docstring = ast.get_docstring(node)
                return docstring if docstring else ''
    except:
        pass
    return ''

def get_docx_unique_workspace_records__1aff24635b2ad5253401338faad2d9bc(env, config):
    """Extract data from docx file keeping only unique workspace IDs.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (e.g., {'path': '/home/user/file.docx'})

    Returns:
        dict: Contains 'workspace_ids' (list), 'is_unique' (bool), 'line_count' (int), 'unique_count' (int)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'workspace_ids': [], 'is_unique': False, 'line_count': 0, 'unique_count': 0}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        workspace_ids = []
        lines = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts = para.text.split(',')
                if len(parts) == 4:
                    workspace_id = parts[3].strip()
                    workspace_ids.append(workspace_id)
                    lines.append(para.text.strip())
        is_unique = len(workspace_ids) == len(set(workspace_ids))
        return {'workspace_ids': workspace_ids, 'is_unique': is_unique, 'line_count': len(lines), 'unique_count': len(set(workspace_ids))}
    except Exception as e:
        logger.error(f'Error processing docx file: {e}')
        return {'workspace_ids': [], 'is_unique': False, 'line_count': 0, 'unique_count': 0}
    finally:
        os.unlink(tmp_path)

def get_vscode_file_content__baf14747(env, config: Dict[str, Any]) -> str:
    """Get the full content of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path'

    Returns:
        The full file content as a string
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        return f.read()

def get_vscode_text_content__681a97c35eb849e5cf9422adfe4e1aea(env, config):
    """Get text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.warning('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_vscode_text_content__65dee18f1880f7d028c2b1727fd62d90(env, config):
    """Extract text content from a file for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the text file, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file {file_path}: {e}')
        return ''

def get_vscode_settings__37a81987aef2f239d369b6635778b9ec(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json content as a dictionary.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dictionary containing the parsed settings.json content, or empty dict on error
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {}
        settings_text = file_bytes.decode('utf-8')
        settings_dict = json.loads(settings_text)
        return settings_dict
    except Exception as e:
        return {}

def get_vscode_extension_dir__42aa293803886f1362c8d3d5e33a2703(env, config: dict):
    """
    Check if VSCode extension directory exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_name' field

    Returns:
        dict: Result with 'exists' boolean and 'path' if found
    """
    extension_name = config.get('extension_name', '')
    if not extension_name:
        return {'exists': False, 'path': None}
    command = f"ls -d ~/.vscode/extensions/{extension_name}* 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('status') == 'success':
        output = result.get('output', '').strip()
        if output and output != '':
            return {'exists': True, 'path': output}
    return {'exists': False, 'path': None}

def get_vscode_workspace_opened__56d366f94894cbfe50f3e560644c1cc7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a workspace file is currently opened in VSCode by checking the accessibility tree.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'workspace_path' key

    Returns:
        Dict with 'vscode_running', 'workspace_opened', 'workspace_name', 'file_exists', 'valid_json', 'folders' keys
    """
    workspace_path = config.get('workspace_path', '')
    workspace_name = os.path.basename(workspace_path)
    result = {'vscode_running': False, 'workspace_opened': False, 'workspace_name': '', 'file_exists': False, 'valid_json': False, 'folders': []}
    file_bytes = env.controller.get_file(workspace_path)
    result['file_exists'] = file_bytes is not None
    if file_bytes:
        try:
            workspace_content = file_bytes.decode('utf-8')
            workspace_data = json.loads(workspace_content)
            result['valid_json'] = True
            result['folders'] = workspace_data.get('folders', [])
        except:
            result['valid_json'] = False
            result['folders'] = []
    try:
        accessibility_tree = env.controller.get_accessibility_tree()
        if accessibility_tree:
            tree_str = str(accessibility_tree).lower()
            if 'visual studio code' in tree_str:
                result['vscode_running'] = True
                workspace_basename = workspace_name.replace('.code-workspace', '')
                if workspace_basename.lower() in tree_str or workspace_name.lower() in tree_str or workspace_path.lower() in tree_str:
                    result['workspace_opened'] = True
                    result['workspace_name'] = workspace_name
    except Exception as e:
        logger.warning(f'Error checking accessibility tree: {e}')
    return result

def get_vscode_settings_autosave__06d1442bba585f1fb8ea50a4956df674(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json file content focusing on autosave settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dict containing the settings.json content, or empty dict if file not found
    """
    file_path = config.get('path', '/home/user/project/.vscode/settings.json')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'Could not read file: {file_path}')
            return {}
        settings = json.loads(file_bytes.decode('utf-8'))
        return settings
    except Exception as e:
        logger.error(f'Error reading VS Code settings: {e}')
        return {}

def get_vscode_file_content__06f73a877cbecd1e2a247622c3136810(env, config):
    """Read text file content from VM for indentation validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_vscode_text_content__62f7a00e7dd5e4db70eb00a615f012ef(env, config):
    """Extract text content from a file for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the text file, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file {file_path}: {e}')
        return ''

def get_vscode_file_content__219e3148d77134036e282c6dd4b41d12(env, config):
    """Read text file content from VM for indentation validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_vscode_file_line__8f9a3936(env, config: Dict[str, Any]) -> str:
    """Get a specific line from a file opened in VSCode.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'line_number'

    Returns:
        The content of the specified line (0-indexed)
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    line_number = config.get('line_number', 0)
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        if line_number < len(lines):
            return lines[line_number].rstrip('\n')
        return ''

def get_vscode_line_count__b2b950a6(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get the total number of lines in a file and its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path'

    Returns:
        Dict with 'line_count' (int) and 'content' (str)
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        content = f.read()
        line_count = len(content.splitlines())
    return {'line_count': line_count, 'content': content}

def get_vscode_search_result__0da38d8e(env, config: Dict[str, Any]) -> str:
    """Get the search results from VS Code.

    This function executes a VS Code extension command to retrieve the actual
    search results (files, matches, line numbers) from a performed search,
    not just the search query text.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file containing search results
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_content__39395653(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_file_content__75a1331b(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_text_content__c5c5b95c23a9d2c28863362320aba24b(env, config):
    """Extract text content from a file for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the text file, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file {file_path}: {e}')
        return ''

def get_vscode_text_content__0f72ff3b9c8d9680b46915659d675f48(env, config):
    """Get text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.warning('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_vscode_variable_usage__f33f7a7e(env, config: Dict[str, Any]) -> int:
    """Count occurrences of a variable name in the file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'variable_name'

    Returns:
        The count of variable occurrences
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    import re
    path = config['path']
    variable_name = config.get('variable_name', '')
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        content = f.read()
        pattern = '\\b' + re.escape(variable_name) + '\\b'
        return len(re.findall(pattern, content))

def get_vscode_open_file__ddd2693d(env, config: Dict[str, Any]) -> str:
    """Get the currently open file in VSCode via the custom extension command."""
    from desktop_env.evaluators.getters.replay import get_replay
    from desktop_env.evaluators.getters.file import get_vm_file
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_open_file__e925f44c(env, config: Dict[str, Any]) -> str:
    """Get the currently open file in VS Code editor.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, dest

    Returns:
        str: Content of the file indicating which file is open
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    from desktop_env.evaluators.getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_line_indentation__4d2e0474(env, config: Dict[str, Any]) -> int:
    """Get the indentation level of a specific line.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'line_number'

    Returns:
        The number of leading spaces
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    line_number = config.get('line_number', 0)
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        if line_number < len(lines):
            line = lines[line_number]
            return len(line) - len(line.lstrip(' '))
        return 0

def get_vscode_file_content__02be85da(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_text_content__dd369016ed084e5c5f565139bb4ef07f(env, config):
    """Extract text content from a file for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the text file, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file {file_path}: {e}')
        return ''

def get_vscode_lines_order__2c3b878a(env, config: Dict[str, Any]) -> str:
    """Get a specific sequence of lines to verify order.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path', 'start_line', and 'end_line'

    Returns:
        Concatenated lines as a single string
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    import os
    path = config['path']
    start_line = config.get('start_line', 0)
    end_line = config.get('end_line', 0)
    file_path = get_vm_file(env, {'path': path, 'dest': os.path.basename(path)})
    with open(file_path, 'r') as f:
        lines = f.readlines()
        if start_line < len(lines) and end_line < len(lines):
            result = []
            for i in range(start_line, end_line + 1):
                if i < len(lines):
                    result.append(lines[i].strip())
            return '|'.join(result)
        return ''

def get_vscode_python_path__1f6f3af3ee1e7e72d7b32984543de005(env, config: Dict[str, Any]) -> str:
    """Get Python path setting from VS Code settings.json.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String containing Python path value, or empty string if not found
    """
    file_path = config.get('path', '/home/user/project/.vscode/settings.json')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'settings.json not found: {file_path}')
            return ''
        settings = json.loads(file_bytes.decode('utf-8'))
        python_path = settings.get('python.defaultInterpreterPath', '')
        return python_path
    except Exception as e:
        logger.error(f'Error reading Python path from settings: {e}')
        return ''

def get_vm_file(env, config: Dict[str, Any]):
    """Get file from VM - simplified version for this getter."""
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        return cache_path
    except Exception as e:
        logger.error(f'Error getting file {path}: {e}')
        return None

def get_vscode_settings_content__a74ebb334721e27ebd0b52511b5b5b07(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get the content of VS Code settings.json file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dict containing the settings, or empty dict if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Could not read file: {config['path']}")
        return {}
    try:
        content = file_bytes.decode('utf-8')
        settings = json.loads(content)
        return settings
    except Exception as e:
        logger.error(f'Error parsing settings.json: {e}')
        return {}

def get_vscode_settings__d0166485dd9dee18fc251c7dbffffc79(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json content as a dictionary.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dictionary containing the parsed settings.json content, or empty dict on error
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {}
        settings_text = file_bytes.decode('utf-8')
        settings_dict = json.loads(settings_text)
        return settings_dict
    except Exception as e:
        return {}

def get_vscode_settings__052d7bff3ecc5844bf94710602cff931(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get VS Code settings.json content as a dictionary.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to settings.json

    Returns:
        Dictionary containing the parsed settings.json content, or empty dict on error
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {}
        settings_text = file_bytes.decode('utf-8')
        settings_dict = json.loads(settings_text)
        return settings_dict
    except Exception as e:
        return {}

def get_vscode_file_content__9d7d44f7(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_explorer__71b2b2791188a1049bb680b7c56ffb2d(env, config: Dict[str, Any]) -> str:
    """Get the VSCode explorer view content using extension command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with vscode_extension_command, path, and dest

    Returns:
        str: Path to the file containing the explorer content
    """
    from ..getters.file import get_vm_file
    from ..getters.replay import get_replay
    os_type = env.vm_platform
    vscode_extension_command = config['vscode_extension_command']
    if os_type == 'MacOS':
        trajectory = [{'type': 'hotkey', 'param': ['command', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    else:
        trajectory = [{'type': 'hotkey', 'param': ['ctrl', 'shift', 'p']}, {'type': 'typewrite', 'param': vscode_extension_command}, {'type': 'press', 'param': 'enter'}]
    get_replay(env, trajectory)
    time.sleep(1.0)
    return get_vm_file(env, {'path': config['path'], 'dest': config['dest']})

def get_vscode_file_content__2fa9f229(env, config):
    """Get file content from VM."""
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, config)

def get_vscode_settings__4df3f2e6(env, config: Dict[str, Any]) -> str:
    """
    Get VS Code settings.json file for format on save checking.

    Args:
        env: Environment instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the downloaded settings file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    return get_vm_file(env, {'path': config.get('path', '/home/user/.config/Code/User/settings.json'), 'dest': config.get('dest', 'settings.json')})

def get_vscode_python_path__e7a35ed0ec20ca7ab6b257f3f5c87e23(env, config: Dict[str, Any]) -> str:
    """Get the Python path setting from VSCode settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for settings file location

    Returns:
        String containing the Python path setting, or empty string if not found
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        settings_json = json.loads(file_bytes.decode('utf-8'))
        return settings_json.get('python.pythonPath', '')
    except Exception as e:
        logger.error(f'Error parsing settings.json: {e}')
        return ''
