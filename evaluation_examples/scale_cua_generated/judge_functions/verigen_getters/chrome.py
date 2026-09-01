"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_url_from_text__3b2adddad1034fb6d584d5542b1c7034', 'get_googledrive_folder_check__e617d7dc', 'get_chrome_sync_disabled__97c301f6', 'get_chrome_session_restore_setting__9cb5a6acff7768f53f96a10b1f86ac95', 'get_writer_table_filtered_rows__ff6673f3b77179ec900fd2f9ae4a2104', 'get_extension_name_path__d51a0e10', 'get_csv_table_biz__fb9ae858', 'get_chrome_clear_cookies_on_exit__37d7b439', 'get_gdrive_file_exists__b791b56367b138c183fb013d8a0662b9', 'get_extension_enabled_state__1e533eef', 'get_cookies_clear_on_exit__207f7b6597ba675431e17149c4a63a1e', 'get_docx_table_data__e6dd308b0cb59f24a41ac5c58e5e6ff5', 'get_gdrive_single_file__54c9c002bd0d5132155806dba06b543e', 'get_csv_table_ee__c4163e13', 'get_popup_html_content__94d6a0b8', 'get_csv_table_cs3y__26538d5a', 'get_chrome_hardware_accel__8f79fa7c', 'get_googledrive_sorted_filenames__939aa00d', 'get_googledrive_eml_list__6866e692', 'get_chrome_cookies_setting__1d96da57', 'get_extension_details__ae6416e4', 'get_extension_names__de268faa9ad661d5adcd42e1b7f775e2', 'get_docx_last_table_dims__3769d65b', 'get_gdrive_pdf_file__f93e2dd2', 'get_gdrive_file_check__a29f27f0', 'get_gdrive_pdf_file__3a6a494e', 'get_third_party_cookies_blocked__ab79f9e8', 'get_chrome_newtab_startup__6cc337d6467b2f92dee1f171f61ea4ed', 'get_gdrive_file_list__59579682', 'get_chrome_bookmarks__b5b6115e', 'get_extension_manifest_version__433004c1c6e2288dd8032cc0aeda6796', 'get_docx_table_content__50486c0ece7eb424890da5a17471074c', 'get_pdf_file_and_chrome_tab__a606126a', 'get_urls_from_txt__a2f501d6', 'get_active_url_from_accessTree', 'get_docx_table_dimensions__76f174dadfebd8762f6007a98d987920', 'get_chrome_homepage__9d2ba16da21ed7a3bc9ce454107b63a8', 'get_gdrive_file_list__7fbebc15', 'get_macys_url_parse__29ba7d7c48cbc7f21a5e114affc6f1be', 'get_webext_manifest__c41914cdb60620ceca4e49be63e9e04c', 'get_docx_table_structure__256ba3c478aefa1c15edcf46f5e2019b', 'get_extension_manifest__986ae666fddb67409d6f0937ecd8b146', 'get_chrome_https_only_mode__bef59d93', 'get_chrome_setting_value__8cfb2b13', 'get_chrome_disable_preload__7df16cad', 'get_docx_table_info__94b525d2', 'get_chrome_experiments_not_contains__f72b4886483f6a090bcc6a28ba49acde', 'get_macys_url_parse__6e5c341c9d3ff8b7569de2db0f8a4fd6', 'get_docx_table_col_count__e15e9162b38246c56c8d7ca12c83c648', 'get_chrome_search_engine__26a050291c58f47430d6ebca143c45ff', 'get_gdrive_file_check__661594c9', 'get_chrome_fixed_font_size__87f268be1c9c07cc266e3e54a412bd8c', 'get_docx_table_content__f06c6ff6344501bcedf4d007c342b6fd', 'get_csv_table_mech__9629fd7e', 'get_docx_table_first_row__f2b7a56f41e0323869ae9541f5efcf40', 'get_docx_table_row2_formatting__6412b9b5', 'get_docx_table_data__7e712d1d6a81a74414ae2e1785ec944e', 'get_chrome_download_location__518a9ee4', 'get_third_party_cookies_blocked__75ed94d8e5e2ff083f81c247646b0b38', 'get_docx_table_data__9453c274a20969d63590381e37c06d85', 'get_table_count_only__db5e4b5e', 'get_chrome_startup_urls__715fc21b1c707dd79fc5ab9b6e4df514', 'get_extension_files_count__014e651b', 'get_chrome_font_size__d996a712', 'get_gdrive_pdf_file__9e5c8fdc', 'get_webext_dir__2ec2a9a8', 'get_extension_manifest_version__b7035c23581ef82d8725cee1b3aa987f', 'get_webext_dir__60d86cd5', 'get_chrome_tabs_and_bookmarks__6d016e82', 'get_gdrive_file_list__fd7bb036', 'get_vim_tabstop_config__76faceaeb09651393a582783b81e5798', 'get_webext_manifest__70dff70b667529be05282f3babd2fe6e', 'get_csv_table_math__bdcc0312', 'get_block_third_party_cookies__980fd38b', 'get_url_text__5525d1c8', 'get_webext_dir__c1aa21b2', 'get_macys_url_parse__7dbecc9b465f893baa757d13deeefa17', 'get_extension_icon_exists__db044e4c', 'get_gdrive_file_check__42b1e128', 'get_csv_table_cs4y__5116177c', 'get_docx_table_row_count__ff8ca82a4f791a2f5346fef371999b2c', 'get_csv_file_and_chrome_tab__4e1aca7b', 'get_googledrive_file_list__dd26081a', 'get_extension_manifest__1edda3bd4444a8fb1554a227ff178017', 'get_default_web_browser__865c6c015a0aa9b8277b205e904948d3', 'get_webext_manifest__d1a4f844666cd6560872c801fdefe60a', 'get_table_position__f4eb9543', 'get_extension_description__2224641eff529090c9cc8ad62127e176', 'get_extension_version__ae6416e4', 'get_csv_table_bio__a4ea4d07', 'get_chrome_policy_page_info__f3b19d1e', 'get_docx_table_info__a4c465e5', 'get_gdrive_text_file__a48d3d2ba069ee77685e4821041681b9', 'get_gdrive_pdf_file__9d82c3b2', 'get_chrome_setting_value__df4ecdef', 'get_chrome_min_font_size__f6a59c3b', 'get_chrome_font_size__caf9b8e1', 'get_gdrive_pdf_file__75da4280', 'get_chrome_min_font_size__904067aa', 'get_gdrive_pdf_info__bb91e7693f2a30704f1d1cc79be73950', 'get_html_exists__7b7b0b2f', 'get_extension_manifest_name__e5eabc9a', 'get_html_file_and_chrome_tab__3cb1a587', 'get_webext_dir__88e21052', 'get_chrome_setting_value__fca61186', 'get_googledrive_eml_files', 'get_vim_hlsearch_config__5438ce42ea45fa77c023ecd730e398a5', 'get_docx_last_table_dims__d15d9273', 'get_vim_tabstop_check__e2e649c5246b836f874ad28b723333bc', 'get_html_result__b5ca523a54cc6ae837c08b60601febd6', 'get_chrome_notifications_blocked__4f8fdc8b', 'get_webext_dir__bba937aa', 'get_extension_name__a4b0d616259699774cd5ee2679c8a96c', 'get_chrome_default_encoding__e3f03b16', 'get_gdrive_twitter_receipt__12ff505f1cc77038855fb34142173dbc', 'get_unpacked_extension_count__0f84311e', 'get_docx_table_info__2341f5ea', 'get_chrome_setting_value__0717aaed', 'get_chrome_setting_value__fd2ee811', 'get_chrome_camera_blocked__d2e0663c', 'get_tabs_and_bookmarks__7a5a7856f1b642a4ade91ca81ca0f263000420251221151547', 'get_docx_table_size__e7b89121046aaa7adf88d3e2fb20c6ab', 'get_googledrive_docx__a60c9fbb05dbe8be7b4812ad58480d11', 'get_csv_table_phys__5e1848a9', 'get_chrome_fixed_font_size__b7f59f23', 'get_html_file_info__48e4da460d6f3e132e4d3cc48ac9ca24', 'get_extension_source_type__ae6416e4', 'get_extension_version__f107c64c583fbfb012ea91827c8f61e3', 'get_chrome_search_engine_changed__6c30ec39', 'get_gdrive_file_list__5876322dac1e30402171f6bcd2edb019', 'get_googledrive_file_list__2b29a90e', 'get_gdrive_file_check__adf6c9b9', 'get_chrome_min_logical_font__b612e026', 'get_chrome_safe_browsing__2b56847e', 'get_chrome_default_search__8a4fe4a4', 'get_extension_icon__0c7a03ab', 'get_extension_popup__8aee113e', 'get_chrome_setting_value__5936c775', 'get_docx_table_structure__8586f38ff35ab9e474999a7f667bbc31', 'get_html_file_and_chrome_tab__d6487d33', 'get_chrome_experiments_exact_match__d7481e87a8a6dde50669ce517c215edf', 'get_block_third_party_cookies__62856a905c85744fe1f63d2f847937c7', 'get_chrome_clear_on_exit__da95206e', 'get_extension_description__221bd7ef80e48e21d62e7d38635cf26c', 'get_recreation_echocanyon_html__dbcda1af4d2b810312479cfb54a15ab6', 'get_chrome_disable_safe_browsing__19371712', 'get_docx_table_info__7c0843e4', 'get_docx_table_content__121c468593688d91f702e8d5fdd1e9f9', 'get_chrome_do_not_track__a858c66a', 'get_chrome_enhanced_safe_browsing__a9ae68b06063b362f29e78385fb296bc', 'get_chrome_location_blocked__d6e72d6c', 'get_table_count_only__d0985f12', 'get_chrome_min_font_enforce__0d902d89d190f0f84e2f19e726b6d1ea', 'get_gdrive_file_metadata__b2463eb9', 'get_extension_count__0fb1bf36ddef8ab5a554de56ff7f0c3d', 'get_googledrive_file_count__17298c22', 'get_docx_table_count__6ea8989b', 'get_show_full_urls__96430b37', 'get_table_position__199b082c', 'get_recreation_url_check__6dc3893f96ccd943c500af1756962de6', 'get_docx_table_row_count__7d2236cd23a11ec2c058cf74d17b7e88', 'get_table_position__6a4e1dd4', 'get_googledrive_has_files__95db6624', 'get_extension_manifest__3d9fc0c33aef9d5f768043f561973d63', 'get_chrome_block_third_party_cookies__327d732b', 'get_chrome_do_not_track__a5d5a0e3', 'get_default_web_browser__b2d61e52', 'get_backup_extension_files__8ece34a1', 'get_recreation_bearlake_html__6adae5a4637b133a46055994fbaa8dd4', 'get_docx_table_structure__98ae47d7', 'get_docx_last_table_dims__356ec724', 'get_chrome_extension_manifest__3d480472b35ce7003ca943ba6b2307fa', 'get_csv_table_civil__2c0c636e', 'get_docx_table_content__a3ef22f0b6dddefc025c6922cfbdc9a0', 'get_gdrive_file_locations__b75faf5b6765d0d1458ed6b6d219047b', 'get_docx_table_content__fefdbfb10ebeb3b7ce54da62c0ba8cd2', 'get_docx_last_table_dims__b98be273', 'get_chrome_enhanced_safe_browsing__a6a55567', 'get_docx_table_header__13eed80defa5785247e65a748bfcefa5', 'get_chrome_extension_manifest__1e834a0e2fee4e317d31e5d8fca95c5c', 'get_table_position__bbe9b961', 'get_extension_manifest__842d772f81b13ef554c1bda7e1c59bb7', 'get_table_position__5ff0bd58', 'get_googledrive_filenames__a2f23245', 'get_csv_table_med__1bf7a84e', 'get_chrome_default_font__ad306c7d', 'get_docx_table_data_rows__969da71b9e6a789c33cf79de761b11d9', 'get_chrome_homepage_setting__32a135b2f85c9e9cdbf4b5150b97474d', 'get_gdrive_file_properties__e9151d35420d1468fe1e721181658d5f', 'get_docx_table_data__6a2751d358b86b079b53effd72604b8f', 'get_table_position__60d42be3', 'get_writer_table_font__e0ebf50c', 'get_extension_dir_exists__79fc6dbb', 'get_gdrive_files_with_pattern__6538c56492a239416edab32324471fe6', 'get_chrome_microphone_blocked__a44d2f59', 'get_file_has_extension__c9bfc81c', 'get_macys_url_parse__972b80b805a61948b132613370427348', 'get_docx_last_table_dims__70149acd', 'get_gdrive_numbered_files__aebf0a91a4be0be45ef3247f943131f2', 'get_chrome_autofill_setting__cb0362f4', 'get_gdrive_file_check__cc11abb5', 'get_extension_has_popup__36b763d0', 'get_gdrive_folder_contents__55cd0f01', 'get_chrome_standard_font_size__c9cdf92f31112c863794561a8e2fdc6b', 'get_chrome_default_font_size__8638fadde1ef14284e07488209d510be', 'get_googledrive_files_with_prefix__fb580d40', 'get_table_row_count_by_index__83874f16', 'get_gdrive_file_metadata__6bd409d3', 'get_chrome_show_home_button__aa9af3d4', 'get_chrome_font_size__82caef9e', 'get_chrome_font_size__fc8c3e93', 'get_gdrive_file_list__da0d777d0a1e3c8d735dbb81a2e45a5c', 'get_gdrive_file_list__b0d15a73736689e26424d81f759f1528', 'get_docx_table_info__818c616c', 'get_default_web_browser__48359cc6c2268bcbb20cc6eebb7a0011', 'get_googledrive_files_list__d6fb1e53c50621e1a08efd7623119b0d', 'get_docx_table_content__044cff297ed58ab26ee69e56f58c3d19', 'get_recreation_cedarbreaks_html__ed948b8fac72e2c98dd5028aabd38bf3', 'get_recreation_table_header__8356d858fc4a70c81a86864f34a14436', 'get_chrome_ext_last_two__f117f7ab324ab80f0f8ad254a42cb210', 'get_recreation_search_box__e0d56e1bf714f2c8f287d0502b986b63', 'get_gdrive_eml_files_with_filenames__74b11cf6', 'get_html_files_in_dir__98346f3b', 'get_chrome_ext_subset__d7f8f8dc6935bc142e9b5dc629fdadfe', 'get_gdrive_file_list__12a331f2', 'get_gdrive_file_list__839d38d5', 'get_table_position__2d408821', 'get_docx_last_table_dims__3c609656', 'get_chrome_download_prompt__e5560b89', 'get_docx_table_count__8c6fecde', 'get_chrome_popups_blocked__6bbb0d83', 'get_table_position__dd5268aa', 'get_writer_table_first_n_rows__92370bc17bf2e2a5bd844c976eb32eb3', 'get_gdrive_simplified_names__43be33c33f047d68d3b6cab5bd4d55ce', 'get_gdrive_file_check__41049d6b', 'get_gdrive_pdf_file__82f685a6', 'get_last_table_structure__e0b81f05', 'get_docx_last_table_dims__f8b20751', 'get_chrome_page_zoom__1c67d2f9', 'get_bookmarks_in_folder__5641e4c0', 'get_csv_file_and_chrome_tab__29a47154', 'get_gdrive_text_file__953256df836603c8857d4495861e4b63', 'get_chrome_font_size__88fb57a4', 'get_chrome_ext_single__abfa4b043befc7603aff32afef71bd1b', 'get_chrome_third_party_cookies_blocked__15369290', 'get_chrome_font_size__0ce3b14e', 'get_recreation_antelope_html__2aa2c046ed13ceae9537c9a8e566d558', 'get_macys_url_parse__78b9ef3c3ae5dc84f1b8b31df8c5818e', 'get_docx_table_info__faf951c7', 'get_chrome_setting_value__a785d845', 'get_chrome_extension_manifest__6e0f538e8ca610f13ae0673161924889', 'get_chrome_disable_fonts__df950818', 'get_gdrive_file_list__87199839', 'get_chrome_min_font_size__03b68be0b364e7ac1421db0dce73b670', 'get_gdrive_pdf_file__bf7218b9', 'get_chrome_ext_middle__484fc06534818b90d35ad638d0805c7c', 'get_chrome_startup_url_removed__a85eebd24e563a97c17935bc46126aa1', 'get_docx_table_structure__aa3c98fadda85fd2b11814a9c969fa8c', 'get_chrome_font_size__25d0e193', 'get_vim_hlsearch_check__47d76eea6c988a4beb0cae6b4109cc24', 'get_gdrive_eml_files__26d2516b888edf7c1b328cca7acaf9b7', 'get_chrome_third_party_cookies__1394774d', 'get_gdrive_filenames__71c23132811d122bc61dca33636b8f81', 'get_docx_table_info__f2fed625', 'get_gdrive_file_list__45753efe', 'get_gdrive_file_check__b2131ee6', 'get_writer_table_rows__b3ed6b27f25c4e432395f7240c942359', 'get_gdrive_pdf_file__c340b25e', 'get_table_count_and_first_table__638d8a63', 'get_docx_table_info__bf2c005c', 'get_docx_table_structure__deda270825a0b396cee34e9436d907d9', 'get_chrome_font_size__baaf192c', 'get_gdrive_file_exists__3ce93f6ce10147dcc2489b34874a9f3b', 'get_webext_dir__951cee96', 'get_chrome_experiments_partial__70b32da51f968d0aa45e836eecc52bdb', 'get_gdrive_file_check__94eb0e23', 'get_chrome_min_font_size__a5307030', 'get_gdrive_file_list__f478ac41', 'get_git_remote_url__f8582c17', 'get_docx_last_table_dims__86f8b2c6', 'get_gdrive_file_check__9b0aebe9', 'get_extension_manifest_version__fb0f6c5b', 'get_chrome_min_font_size__986d3358722e45a6a9be7a39f21f689e', 'get_googledrive_files_list__ba61b137046435f47239f8911466a875', 'get_extension_manifest__b5fa3477caae1bab3fd0d8b19ef4dfc7', 'get_gdrive_nested_files__688ade84b6705b21f2a27dc4863ba216', 'get_chrome_fixed_font_family__a0c1ef41', 'get_chrome_password_breach_alerts__2107d226', 'get_clear_cookies_on_exit__936eec5ac877849802aeaa5f0c43e44f', 'get_webext_dir__81e9cd46', 'get_docx_table_content__0da226e8b4b1ba563f27816711980e35', 'get_chrome_zoom_level__d32e355a', 'get_webext_manifest__4b1e45b68c5d85573f69177601102dcb', 'get_chrome_ext_productivity__3ce2c199effd9eefc89344b99f4cd5a7', 'get_docx_table_data__54b06cb187414c0c4afc1326d01127a0', 'get_docx_table_content__1c032efed907d41d75501f3895adca32', 'get_chrome_password_manager__09e00537', 'get_gdrive_file_check__ada1e84e', 'get_gdrive_pdf_file__7e0bdf48', 'get_extensions_page_content__ae6416e4', 'get_extension_manifest_version__ae6416e4', 'get_docx_table_bold_status__67405416', 'get_gdrive_pdf_file__a70aba9d', 'get_docx_table_info__2aa481b8', 'get_chrome_extension_manifest__70453fbf044a41274b6b7dc7e909fb11', 'get_chrome_extension_manifest__61167ddb747c3a88a31446374f5a958f', 'get_block_third_party_cookies__d0393d13df595b6db99860dc4f30ea7b', 'get_googledrive_folder_details__68f5fe6b', 'get_chrome_font_size__4edc5823', 'get_gdrive_text_file__a652858db2e92fae77817389157c8edc', 'get_table_position__b943f0fd', 'get_docx_last_table_dims__cbb3bae6', 'get_docx_table_data__d5cdfdb611de0d2892eeaddda9638ac9', 'get_table_position__d963ccf8', 'get_chrome_setting_value__eec6beed', 'get_googledrive_folder_file_count__1271f790', 'get_extension_version__a366045b', 'get_docx_table_borders__5aa272b32299efe042e474c4fb5400ce', 'get_chrome_setting_value__d84aae11', 'get_docx_table_content__0918fd6ddecb95fb671b970e39bbe53c', 'get_docx_table_data__0c283e204015cae1eab9b0792877891d', 'get_docx_last_table_dims__c936c1ec', 'get_open_tab_count__a852a89d', 'get_extension_enabled_state__6f71517e0c42749af1a6363d9f36e224', 'get_chrome_location_blocked__6591cb23', 'get_googledrive_backup_check__c15ab8b3', 'get_extension_description__db4e5321', 'get_docx_arxiv_urls__82d38938', 'get_extension_enabled_status__ae6416e4', 'get_chrome_default_font_size__e289225d', 'get_default_web_browser__88e8e17a', 'get_html_file_and_chrome_tab__6920e35a', 'get_chrome_setting_value__1be3beb2', 'get_all_extension_paths__225a261e', 'get_docx_table_info__fca84b62', 'get_chrome_javascript_setting__da00e3f2', 'get_docx_table_data__b0edb7cf3cae7467b1751eb74a239c0d', 'get_writer_table_with_headers__cb84d8e55491bfd63123932189f4803b', 'get_extension_description__ae6416e4', 'get_table_count_only__e164fbd9', 'get_chrome_sansserif_font__2cf2a146', 'get_docx_table_content__28c7c94bbd8c2d75a3d0fb78ff321239', 'get_chrome_experiments_contains__0a40109ea287ab7bbd8cd9175a7ce6a5', 'get_extension_folder_exists__ae6416e4', 'get_gdrive_file_check__db0d2d11', 'get_chrome_serif_font__ddb77dce', 'get_webext_manifest__dea4a5ea3d588b414bf151fee72b35b0', 'get_chrome_experiments_multi__e998f78abb27064086318477b860256b', 'get_chrome_fixed_font_size__d608837f75cec39dc6023178df898af7', 'get_webext_dir__b9b28e4a', 'get_sorted_table_data__857e1276', 'get_recreation_devilsgarden_html__fa1e76c31141f93d38de11c4bb8239cf', 'get_file_extension__e13be972', 'get_googledrive_folder_info__a18b8359', 'get_extension_version__407be0458b7b234fb5401d66a10f5221', 'get_gdrive_file_list__3dcb78c90ac64690f9f399090d59db08', 'get_chrome_font_size__9b3236ac']

def get_url_from_text__3b2adddad1034fb6d584d5542b1c7034(env, config):
    """
    Extract a URL from a text file.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key specifying the file path on VM

    Returns:
        str: The URL found in the file, or empty string if not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes:
        try:
            import re
            content = file_bytes.decode('utf-8', errors='ignore')
            match = re.search('https?://[^\\s]+', content)
            if match:
                return match.group(0).strip()
            return ''
        except Exception:
            return ''
    else:
        return ''

def get_googledrive_folder_check__e617d7dc(env, config: Dict[str, Any]):
    """Check if folder exists and count files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: {"folder_exists": bool, "file_count": int}
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_name = config.get('folder_name', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and 'root' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        if len(filelist) == 0:
            return {'folder_exists': False, 'file_count': 0}
        folder_id = filelist[0]['id']
        file_query = f"'{folder_id}' in parents and trashed = false"
        file_list = drive.ListFile({'q': file_query}).GetList()
        file_count = sum((1 for f in file_list if f['mimeType'] != 'application/vnd.google-apps.folder'))
        return {'folder_exists': True, 'file_count': file_count}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'folder_exists': False, 'file_count': 0}

def get_chrome_sync_disabled__97c301f6(env, config: Dict[str, str]):
    """
    Check if Chrome Sync is disabled.
    Returns 'true' if sync is disabled, 'false' if enabled.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        sync_requested = data.get('sync', {}).get('requested', False)
        signin = data.get('signin', {})
        signin_allowed = signin.get('allowed', True)
        sync_disabled = not sync_requested and (not signin_allowed)
        explicit_disabled = data.get('sync_disabled', False)
        return 'true' if sync_disabled or explicit_disabled or (not signin_allowed) else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome sync status: {e}')
        return 'false'

def get_chrome_session_restore_setting__9cb5a6acff7768f53f96a10b1f86ac95(env, config: Dict[str, str]):
    """
    Get the Chrome session restore setting.
    Returns the restore_on_startup value (1-5).
    1 = Reopen last session
    4 = Open a specific page or set of pages
    5 = Open the New Tab page
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        restore_setting = data.get('session', {}).get('restore_on_startup', 5)
        logger.info(f'Current restore_on_startup: {restore_setting}')
        return restore_setting
    except Exception as e:
        logger.error(f'Error getting restore setting: {e}')
        return 5

def get_writer_table_filtered_rows__ff6673f3b77179ec900fd2f9ae4a2104(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract filtered rows from Writer document table.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with table data
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            if len(doc.tables) == 0:
                return {'error': 'No tables found'}
            table = doc.tables[0]
            rows_data = []
            for row in table.rows:
                row_values = []
                for cell in row.cells:
                    text = cell.text.strip()
                    try:
                        if '.' in text:
                            val = float(text)
                        elif text.isdigit() or (text and text[0].isdigit()):
                            val = int(text)
                        else:
                            val = text
                        row_values.append(val)
                    except:
                        row_values.append(text)
                rows_data.append(row_values)
            return {'row_count': len(rows_data), 'rows': rows_data, 'has_header': len(rows_data) > 0}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_extension_name_path__d51a0e10(env, config):
    """
    Get all installed Chrome extensions with their names and paths.

    Returns a list of dicts containing extension information:
    [
        {
            "name": "Extension Name",
            "path": "/path/to/extension"
        },
        ...
    ]
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        logger.error(f'Unsupported operating system: {os_type}')
        return []
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extensions_info = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for (extension_id, extension_data) in all_extensions.items():
            name = extension_data.get('manifest', {}).get('name', '')
            path = extension_data.get('path', '')
            if name and path:
                extensions_info.append({'name': name, 'path': path})
                logger.info(f"Found extension: name='{name}', path='{path}'")
        logger.info(f'Total extensions found: {len(extensions_info)}')
        return extensions_info
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return []

def get_csv_table_biz__fb9ae858(env, config):
    """Get CSV table content for biz task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/BIZ-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_chrome_clear_cookies_on_exit__37d7b439(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to clear cookies on exit.
    Returns 'true' if setting is enabled, 'false' otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        content_settings = profile.get('content_settings', {})
        exceptions = content_settings.get('exceptions', {})
        cookies_settings = exceptions.get('cookies', {})
        exit_type = profile.get('exit_type', '')
        clear_on_exit = profile.get('clear_cookies_on_exit', False)
        default_content_settings = profile.get('default_content_setting_values', {})
        cookies_session_only = default_content_settings.get('cookies', 0) == 4
        return 'true' if clear_on_exit or cookies_session_only else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome clear cookies on exit setting: {e}')
        return 'false'

def get_gdrive_file_exists__b791b56367b138c183fb013d8a0662b9(env, config: Dict[str, Any]) -> bool:
    """Check if a file exists in a specific Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with Google Drive settings and path

    Returns:
        bool: True if file exists in the specified folder, False otherwise
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
    except ImportError as e:
        logger.error(f'Missing required library: {e}')
        return False
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    path_list = config.get('path', [])
    if not path_list:
        logger.warning('No path specified in config')
        return False
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for (idx, folder_or_file) in enumerate(path_list):
            is_folder = idx < len(path_list) - 1
            if is_folder:
                search = f"title = '{folder_or_file}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false and '{parent_id}' in parents"
            else:
                search = f"title = '{folder_or_file}' and trashed = false and '{parent_id}' in parents"
            filelist = drive.ListFile({'q': search}).GetList()
            if len(filelist) == 0:
                logger.info(f"File/folder '{folder_or_file}' not found in Google Drive at level {idx}")
                return False
            file = filelist[0]
            parent_id = file['id']
        logger.info(f"File found in Google Drive: {'/'.join(path_list)}")
        return True
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return False

def get_extension_enabled_state__1e533eef(env, config):
    """
    Get the enabled state of a Chrome extension by name.

    This function reads the Chrome Preferences file to check if an extension
    with the given name is both installed and enabled (state=1).

    Args:
        env: Environment object
        config: Configuration dict containing:
            - extension_name (str): The name of the extension to check

    Returns:
        dict: Dictionary containing:
            - extension_name (str): The name of the extension
            - is_installed (bool): Whether the extension is installed
            - is_enabled (bool): Whether the extension is enabled (state=1)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extension_name = config.get('extension_name', '')
        logger.info(f'[EXTENSION_ENABLED_STATE] Looking for extension: {extension_name}')
        all_extensions = data.get('extensions', {}).get('settings', {})
        extension_found = False
        is_enabled = False
        for (extension_id, extension_data) in all_extensions.items():
            manifest = extension_data.get('manifest', {})
            ext_name = manifest.get('name', '')
            logger.info(f'[EXTENSION_ENABLED_STATE] Checking extension ID {extension_id}: {ext_name}')
            if ext_name == extension_name:
                extension_found = True
                state = extension_data.get('state', 0)
                is_enabled = state == 1
                logger.info(f"[EXTENSION_ENABLED_STATE] Found extension '{extension_name}'")
                logger.info(f'[EXTENSION_ENABLED_STATE] State: {state} (1=enabled, 0=disabled)')
                logger.info(f'[EXTENSION_ENABLED_STATE] Is enabled: {is_enabled}')
                break
        if not extension_found:
            logger.info(f"[EXTENSION_ENABLED_STATE] Extension '{extension_name}' not found in Chrome preferences")
        return {'extension_name': extension_name, 'is_installed': extension_found, 'is_enabled': is_enabled}
    except Exception as e:
        logger.error(f'[EXTENSION_ENABLED_STATE] Error reading Chrome Preferences file: {e}')
        return {'extension_name': config.get('extension_name', ''), 'is_installed': False, 'is_enabled': False}

def get_cookies_clear_on_exit__207f7b6597ba675431e17149c4a63a1e(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to clear cookies when the browser closes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: "true" if cookies are set to clear on exit, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile_data = data.get('profile', {})
        default_content_settings = profile_data.get('default_content_setting_values', {})
        cookies_default = default_content_settings.get('cookies', 1)
        if cookies_default == 4:
            return 'true'
        content_settings = profile_data.get('content_settings', {})
        exceptions = content_settings.get('exceptions', {})
        cookies_exceptions = exceptions.get('cookies', {})
        for (pattern, pattern_data) in cookies_exceptions.items():
            if isinstance(pattern_data, dict):
                for (origin, origin_data) in pattern_data.items():
                    if isinstance(origin_data, dict):
                        setting_value = origin_data.get('setting', 1)
                        if setting_value == 4:
                            return 'true'
        return 'false'
    except Exception as e:
        logger.error(f'Error checking cookies clear on exit: {e}')
        return 'false'

def get_docx_table_data__e6dd308b0cb59f24a41ac5c58e5e6ff5(env, config):
    """Extract table data from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists containing table cell values
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return []
        table = doc.tables[0]
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                try:
                    if '.' in cell_text:
                        row_data.append(float(cell_text))
                    else:
                        row_data.append(int(cell_text))
                except (ValueError, AttributeError):
                    row_data.append(cell_text)
            table_data.append(row_data)
        return table_data
    finally:
        os.unlink(tmp_path)

def get_gdrive_single_file__54c9c002bd0d5132155806dba06b543e(env, config: Dict[str, Any]) -> Optional[str]:
    """Get a single email file from Google Drive by exact filename.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_name: Name of the folder containing file (default: 'emails')
            - filename: Exact filename to search for
            - dest: Local destination filename

    Returns:
        Local filepath where file was downloaded, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    filename = config.get('filename')
    dest = config.get('dest', 'pred.eml')
    if not filename:
        logger.error('No filename specified in config')
        return None
    auth = GoogleAuth(settings_file=settings_file)
    drive = GoogleDrive(auth)
    folder_query = f"title = '{folder_name}' and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
    folder_list: GoogleDriveFileList = drive.ListFile({'q': folder_query}).GetList()
    if len(folder_list) == 0:
        logger.warning(f"Folder '{folder_name}' not found in Google Drive")
        return None
    folder_id = folder_list[0]['id']
    file_query = f"title = '{filename}' and '{folder_id}' in parents"
    file_list: GoogleDriveFileList = drive.ListFile({'q': file_query}).GetList()
    if len(file_list) == 0:
        logger.warning(f"File '{filename}' not found in folder '{folder_name}'")
        return None
    file: GoogleDriveFile = file_list[0]
    try:
        file.GetContentFile(dest, mimetype=file['mimeType'])
        return dest
    except Exception as e:
        logger.error(f"Failed to download '{filename}': {e}")
        return None

def get_csv_table_ee__c4163e13(env, config):
    """Get CSV table content for ee task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/EE-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_popup_html_content__94d6a0b8(env, config: Dict[str, Any]):
    """Read the popup HTML file content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: HTML content, or None if not found
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    html_path = f'{extension_path}/hello.html'
    try:
        content = env.controller.get_file(html_path)
        if content:
            return content.decode('utf-8') if isinstance(content, bytes) else content
        return None
    except Exception as e:
        logger.error(f'Error reading popup HTML: {e}')
        return None

def get_csv_table_cs3y__26538d5a(env, config):
    """Get CSV table content for cs3y task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/CS-p3y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_chrome_hardware_accel__8f79fa7c(env, config: Dict[str, str]):
    """Get Chrome hardware acceleration status from preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Hardware acceleration settings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        hardware_accel_disabled = data.get('hardware_acceleration_mode', {}).get('enabled', None)
        return {'enabled': hardware_accel_disabled is not False}
    except Exception as e:
        logger.error(f'Error getting hardware acceleration status: {e}')
        return {'enabled': True}

def get_googledrive_sorted_filenames__939aa00d(env, config: Dict[str, Any]) -> List[str]:
    """
    Get sorted list of filenames in Google Drive folder.

    Connects to Google Drive using PyDrive, finds the folder matching the query,
    and returns a sorted list of all non-trashed file titles in that folder.

    Args:
        env: Environment object (unused but required for signature)
        config: Dict with 'settings_file' (path to Google Drive settings YAML)
                and 'folder_query' (Drive API query to find the target folder)

    Returns:
        List of filenames (sorted alphabetically), or empty list if folder not found or error
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            return []
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        filenames = sorted([f['title'] for f in file_list])
        logger.info(f'Files: {filenames}')
        return filenames
    except Exception as e:
        logger.error(f'Error: {e}')
        return []

def get_googledrive_eml_list__6866e692(env, config: Dict[str, Any]):
    """Get list of .eml files in folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: List of .eml filenames
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                return []
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        eml_files = [f['title'] for f in filelist if f.get('title', '').endswith('.eml')]
        return eml_files
    except Exception as e:
        logger.error(f'Error: {e}')
        return []

def get_chrome_cookies_setting__1d96da57(env, config: Dict[str, str]):
    """
    Get the cookies content setting from Chrome preferences.

    Args:
        env: Desktop environment instance
        config: Configuration dictionary

    Returns:
        str: "allow" if cookies are enabled, "block" if disabled
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        cookies_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('cookies', 1)
        return 'block' if cookies_setting == 2 else 'allow'
    except Exception as e:
        logger.error(f'Error getting cookies setting: {e}')
        return 'allow'

def get_extension_details__ae6416e4(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Get detailed information for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        List of dicts, each containing extension details (name, version, description, etc.)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extensions_details = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            manifest = ext_data.get('manifest', {})
            details = {'name': manifest.get('name', ''), 'version': manifest.get('version', ''), 'description': manifest.get('description', ''), 'manifest_version': manifest.get('manifest_version', 0)}
            extensions_details.append(details)
        return extensions_details
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file for extension details: {e}')
        return []

def get_extension_names__de268faa9ad661d5adcd42e1b7f775e2(env, config):
    """Get list of all installed extension names from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of extension names (strings)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extension_names = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            if 'manifest' in ext_data:
                name = ext_data['manifest'].get('name', '')
                if name:
                    extension_names.append(name)
            elif 'path' in ext_data:
                path = ext_data['path']
                import os
                folder_name = os.path.basename(path)
                if folder_name:
                    extension_names.append(folder_name)
        return extension_names
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return []

def get_docx_last_table_dims__3769d65b(env, config: Dict[str, Any]):
    """Get table count, dimensions and content status of the last table in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'table_count', 'rows', 'cols', and 'is_blank' keys, or None if error
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        table_count = len(doc.tables)
        if table_count == 0:
            return None
        last_table = doc.tables[-1]
        num_rows = len(last_table.rows)
        num_cols = len(last_table.columns)
        is_blank = True
        for row in last_table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    is_blank = False
                    break
            if not is_blank:
                break
        return {'table_count': table_count, 'rows': num_rows, 'cols': num_cols, 'is_blank': is_blank}
    except Exception as e:
        return None

def get_gdrive_pdf_file__f93e2dd2(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_gdrive_file_check__a29f27f0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_gdrive_pdf_file__3a6a494e(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_third_party_cookies_blocked__ab79f9e8(env, config: Dict[str, str]):
    """Check if third-party cookies are blocked in Chrome settings."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        block_third_party_cookies = data.get('profile', {}).get('block_third_party_cookies', False)
        return 'true' if block_third_party_cookies else 'false'
    except Exception as e:
        logger.error(f'Error: {e}')
        return 'false'

def get_chrome_newtab_startup__6cc337d6467b2f92dee1f171f61ea4ed(env, config: Dict[str, str]):
    """
    Check if Chrome is set to open New Tab page on startup.
    Returns "true" if set to New Tab page, "false" otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        if 'session' not in data.keys():
            return 'true'
        else:
            restore_setting = data.get('session', {}).get('restore_on_startup', 5)
            logger.info(f'Current restore_on_startup: {restore_setting}')
            return 'true' if restore_setting == 5 else 'false'
    except Exception as e:
        logger.error(f'Error checking New Tab startup: {e}')
        return 'false'

def get_gdrive_file_list__59579682(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_chrome_bookmarks__b5b6115e(env, config: dict):
    """Get Chrome bookmarks.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: List of bookmarks
    """
    try:
        bookmarks = get_bookmarks(env, {})
        if bookmarks:
            logger.info(f'Found {len(bookmarks)} bookmarks')
            return bookmarks
    except Exception as e:
        logger.error(f'Error getting bookmarks: {e}')
    return []

def get_extension_manifest_version__433004c1c6e2288dd8032cc0aeda6796(env, config):
    """Get manifest_version of a specific extension by name.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_name' key

    Returns:
        int: Manifest version of the extension, or 0 if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    target_name = config.get('extension_name', '')
    if not target_name:
        return 0
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            if 'manifest' in ext_data:
                name = ext_data['manifest'].get('name', '')
                if name.lower() == target_name.lower():
                    manifest_version = ext_data['manifest'].get('manifest_version', 0)
                    return manifest_version
        return 0
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return 0

def get_docx_table_content__50486c0ece7eb424890da5a17471074c(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract table content from a DOCX file with section context.

    This getter reads a DOCX file and extracts all tables, returning their content
    as a structured dictionary with information about the section they appear in.

    Args:
        env: DesktopEnv instance with controller
        config: Configuration dict with 'path' key pointing to the DOCX file on VM

    Returns:
        Dictionary with tables data and section context, or None if error occurs
        Format: {
            'num_tables': int,
            'tables': [
                {
                    'num_rows': int,
                    'num_cols': int,
                    'data': [[cell_text, ...], ...],
                    'section_before': str (heading text found before this table)
                },
                ...
            ]
        }
    """
    try:
        file_path = config.get('path', '')
        if not file_path:
            logger.error('No path specified in config')
            return None
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {file_path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            tables_data = []
            current_section = ''
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    para_text = ''
                    for text_elem in element.iter():
                        if text_elem.text:
                            para_text += text_elem.text
                    para_text = para_text.strip()
                    para_lower = para_text.lower()
                    if any((keyword in para_lower for keyword in ['result', 'section', 'chapter', 'introduction', 'conclusion', 'method'])):
                        current_section = para_text
                elif element.tag.endswith('tbl'):
                    for table in doc.tables:
                        if table._element == element:
                            table_info = {'num_rows': len(table.rows), 'num_cols': len(table.columns) if table.rows else 0, 'data': [], 'section_before': current_section}
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_info['data'].append(row_data)
                            tables_data.append(table_info)
                            break
            result = {'num_tables': len(tables_data), 'tables': tables_data}
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting DOCX table content: {e}')
        return None

def get_pdf_file_and_chrome_tab__a606126a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if PDF file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    pdf_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.pdf'
    result = env.controller.get_file(pdf_path)
    file_exists = result is not None and len(result) > 0
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    tabs_info = get_open_tabs_info(env, {})
    chrome_tabs = []
    if tabs_info:
        chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': pdf_path}

def get_urls_from_txt__a2f501d6(env, config):
    """
    Extract URLs from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of URLs found (as strings)
    """
    import re
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return []
        content = file_bytes.decode('utf-8', errors='ignore')
        url_pattern = 'https?://[^\\s<>"{}|\\\\^`\\[\\]]+'
        urls = re.findall(url_pattern, content)
        return urls
    except Exception as e:
        return []

def get_active_url_from_accessTree(env, config: Dict[str, str]):
    """Get active URL from accessibility tree"""
    import lxml.etree
    _accessibility_ns_map = {'st': 'uri:deskat:state.at-spi.gnome.org', 'attr': 'uri:deskat:attributes.at-spi.gnome.org', 'cp': 'uri:deskat:component.at-spi.gnome.org', 'doc': 'uri:deskat:document.at-spi.gnome.org', 'docattr': 'uri:deskat:attributes.document.at-spi.gnome.org', 'txt': 'uri:deskat:text.at-spi.gnome.org', 'val': 'uri:deskat:value.at-spi.gnome.org', 'act': 'uri:deskat:action.at-spi.gnome.org'}
    xml_str = env.controller.get_accessibility_tree()
    if xml_str is None:
        return None
    root = lxml.etree.fromstring(xml_str.encode('utf-8'))
    prefix = config.get('goto_prefix', 'https://www.')
    for frame in root.xpath("//frame[@st:focused='true']", namespaces=_accessibility_ns_map):
        name = frame.get('name')
        if name and name.startswith(prefix):
            return name
    return None

def get_docx_table_dimensions__76f174dadfebd8762f6007a98d987920(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get dimensions (rows and columns) from tables in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        Dict with table information: {'tables': [{'rows': int, 'columns': int}, ...]}
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return {'tables': []}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'tables': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        tables_info = []
        for table in doc.tables:
            tables_info.append({'rows': len(table.rows), 'columns': len(table.columns)})
        return {'tables': tables_info}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_chrome_homepage__9d2ba16da21ed7a3bc9ce454107b63a8(env, config: Dict[str, str]):
    """
    Get the homepage URL from Chrome browser.
    Returns the homepage URL set in Chrome preferences.
    Checks both startup_urls and homepage button field.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    expected_url = config.get('expected', '').lower().strip()
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        homepage = data.get('homepage', None)
        if homepage:
            homepage_normalized = homepage.lower().strip()
            if expected_url in homepage_normalized or homepage_normalized in expected_url:
                return homepage
        startup_urls = data.get('session', {}).get('startup_urls', [])
        for url in startup_urls:
            url_normalized = url.lower().strip()
            if expected_url in url_normalized or url_normalized in expected_url:
                return url
        if startup_urls:
            return startup_urls[0]
        return homepage
    except Exception as e:
        logger.error(f'Error getting Chrome homepage: {e}')
        return None

def get_gdrive_file_list__7fbebc15(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_macys_url_parse__29ba7d7c48cbc7f21a5e114affc6f1be(env, config: Dict[str, str]):
    """
    Parse Macy's product URL for men's medium-size short-sleeve shirts with any discount.
    Variation 4: men's medium-size short-sleeve shirts with any discount.
    """
    result = {}
    active_tab_url = get_active_url_from_accessTree(env, config)
    if active_tab_url is None:
        return None
    parsed = urlparse(active_tab_url)
    path = unquote(parsed.path)
    result['mens_clothing'] = True if 'mens-clothing' in path else None
    path_parts = path.strip('/').split('/')
    key_value_json = {}
    shirts_flag = False
    short_sleeve_flag = False
    long_sleeve_flag = False
    if 'shirts' in path:
        shirts_flag = True
    if 'short-sleeve' in path:
        short_sleeve_flag = True
    if 'long-sleeve' in path:
        long_sleeve_flag = True
    for i in range(len(path_parts) - 1):
        if ',' in path_parts[i] and ',' in path_parts[i + 1]:
            keys = [k.strip() for k in path_parts[i].split(',')]
            values = [v.strip() for v in path_parts[i + 1].split(',')]
            for (k, v) in zip(keys, values):
                if k == 'Price_discount_range':
                    key_value_json[k] = [item.strip() for item in v.split('|')] if v else None
                else:
                    key_value_json[k] = v if v else None
                if k == 'Product_department' and (v == 'shirts' or v == 'Shirts' or v == 'Shirt'):
                    shirts_flag = True
                if k == 'Sleeve_length':
                    if v == 'short-sleeve' or v == 'Short Sleeve':
                        short_sleeve_flag = True
                    elif v == 'long-sleeve' or v == 'Long Sleeve':
                        long_sleeve_flag = True
            break
    for field in ['Men_regular_size_t', 'Price_discount_range', 'Sleeve_length']:
        if field not in key_value_json:
            key_value_json[field] = None
    result['shirts'] = shirts_flag if shirts_flag else None
    result['short_sleeve'] = short_sleeve_flag if short_sleeve_flag else None
    result['long_sleeve'] = long_sleeve_flag if long_sleeve_flag else None
    for key in config.get('parse_keys', []):
        if key in key_value_json:
            if key == 'Price_discount_range':
                if key_value_json[key] is not None:
                    if '50_PERCENT_ off & more' in key_value_json[key] and (not '30_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '50_PERCENT_ off & more'
                    elif '30_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '30_PERCENT_ off & more'
                    elif '20_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '30_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '20_PERCENT_ off & more'
                    else:
                        result[key] = 'other_discount'
                else:
                    result[key] = 'no_discount'
            else:
                result[key] = key_value_json[key]
    return result

def get_webext_manifest__c41914cdb60620ceca4e49be63e9e04c(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract manifest.json from a web extension project directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to manifest.json on VM)

    Returns:
        Dict containing the manifest JSON, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest file from VM: {path}')
            return None
        manifest = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded manifest from {path}')
        return manifest
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading manifest from {path}: {e}')
        return None

def get_docx_table_structure__256ba3c478aefa1c15edcf46f5e2019b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract table structure and key cells from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'table_index'

    Returns:
        Dict containing:
            - exists: bool (whether table exists)
            - row_count: int (number of rows)
            - col_count: int (number of columns)
            - cells: Dict[str, str] (selected cell contents by position key)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        table_index = config.get('table_index', 0)
        if len(doc.tables) <= table_index:
            return {'exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}
        table = doc.tables[table_index]
        row_count = len(table.rows)
        col_count = len(table.columns) if row_count > 0 else 0
        cells = {}
        for row_idx in range(row_count):
            for col_idx in range(col_count):
                try:
                    cell_key = f'r{row_idx}c{col_idx}'
                    cells[cell_key] = table.cell(row_idx, col_idx).text.strip()
                except:
                    pass
        return {'exists': True, 'row_count': row_count, 'col_count': col_count, 'cells': cells}
    except Exception as e:
        return {'exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}
    finally:
        os.unlink(tmp_path)

def get_extension_manifest__986ae666fddb67409d6f0937ecd8b146(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract manifest.json content from browser extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to manifest.json

    Returns:
        Dictionary containing manifest.json content, or empty dict if file not found
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest.json from VM: {path}')
            return {}
        manifest_data = json.loads(file_bytes.decode('utf-8'))
        return manifest_data
    except json.JSONDecodeError as e:
        logger.error(f'JSON decode error in manifest.json: {e}')
        return {}
    except Exception as e:
        logger.error(f'Error reading manifest.json from {path}: {e}')
        return {}

def get_chrome_https_only_mode__bef59d93(env, config: Dict[str, str]):
    """
    Check if Chrome HTTPS-First mode is enabled.
    Returns 'true' if HTTPS-only/HTTPS-first mode is enabled, 'false' otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        generated = data.get('generated', {})
        https_first_enabled = generated.get('https_first_mode_enabled', False)
        profile = data.get('profile', {})
        https_upgrade = profile.get('https_only_mode_enabled', False)
        return 'true' if https_first_enabled or https_upgrade else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome HTTPS-only mode: {e}')
        return 'false'

def get_chrome_setting_value__8cfb2b13(env, config: Dict[str, str]):
    """
    Get the Safe Browsing setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['safebrowsing', 'enhanced']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, False)
            else:
                setting_value = False
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': False}

def get_chrome_disable_preload__7df16cad(env, config: Dict[str, str]):
    """
    Check if Chrome preloading is disabled (for privacy/bandwidth).
    Returns 'true' if preloading is disabled, 'false' if enabled.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        net_prefs = data.get('net', {})
        network_prediction_options = net_prefs.get('network_prediction_options', 2)
        profile = data.get('profile', {})
        profile_prediction = profile.get('network_prediction_options', 2)
        is_disabled = network_prediction_options == 0 or profile_prediction == 0
        return 'true' if is_disabled else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome preload settings: {e}')
        return 'false'

def get_docx_table_info__94b525d2(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_chrome_experiments_not_contains__f72b4886483f6a090bcc6a28ba49acde(env, config: Dict[str, str]):
    """
    Get enabled Chrome experiments and return them as a list.
    This getter is used to check if specific experiments are NOT enabled.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, but required by framework)

    Returns:
        List of enabled experiment names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Local State'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enabled_labs_experiments = data.get('browser', {}).get('enabled_labs_experiments', [])
        experiment_names = [exp.split('@')[0] for exp in enabled_labs_experiments]
        return experiment_names
    except Exception as e:
        logger.error(f'Error getting enabled experiments: {e}')
        return []

def get_macys_url_parse__6e5c341c9d3ff8b7569de2db0f8a4fd6(env, config: Dict[str, str]):
    """
    Parse Macy's product URL for men's large-size shirts (any sleeve) with 50% discount.
    Variation 3: men's large-size shirts (any sleeve) with 50% discount.
    """
    result = {}
    active_tab_url = get_active_url_from_accessTree(env, config)
    if active_tab_url is None:
        return None
    parsed = urlparse(active_tab_url)
    path = unquote(parsed.path)
    result['mens_clothing'] = True if 'mens-clothing' in path else None
    path_parts = path.strip('/').split('/')
    key_value_json = {}
    shirts_flag = False
    short_sleeve_flag = False
    long_sleeve_flag = False
    if 'shirts' in path:
        shirts_flag = True
    if 'short-sleeve' in path:
        short_sleeve_flag = True
    if 'long-sleeve' in path:
        long_sleeve_flag = True
    for i in range(len(path_parts) - 1):
        if ',' in path_parts[i] and ',' in path_parts[i + 1]:
            keys = [k.strip() for k in path_parts[i].split(',')]
            values = [v.strip() for v in path_parts[i + 1].split(',')]
            for (k, v) in zip(keys, values):
                if k == 'Price_discount_range':
                    key_value_json[k] = [item.strip() for item in v.split('|')] if v else None
                else:
                    key_value_json[k] = v if v else None
                if k == 'Product_department' and (v == 'shirts' or v == 'Shirts' or v == 'Shirt'):
                    shirts_flag = True
                if k == 'Sleeve_length':
                    if v == 'short-sleeve' or v == 'Short Sleeve':
                        short_sleeve_flag = True
                    elif v == 'long-sleeve' or v == 'Long Sleeve':
                        long_sleeve_flag = True
            break
    for field in ['Men_regular_size_t', 'Price_discount_range', 'Sleeve_length']:
        if field not in key_value_json:
            key_value_json[field] = None
    result['shirts'] = shirts_flag if shirts_flag else None
    result['short_sleeve'] = short_sleeve_flag if short_sleeve_flag else None
    result['long_sleeve'] = long_sleeve_flag if long_sleeve_flag else None
    for key in config.get('parse_keys', []):
        if key in key_value_json:
            if key == 'Price_discount_range':
                if key_value_json[key] is not None:
                    if '50_PERCENT_ off & more' in key_value_json[key] and (not '30_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '50_PERCENT_ off & more'
                    elif '30_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '30_PERCENT_ off & more'
                    elif '20_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '30_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '20_PERCENT_ off & more'
                    else:
                        result[key] = 'other_discount'
                else:
                    result[key] = 'no_discount'
            else:
                result[key] = key_value_json[key]
    return result

def get_docx_table_col_count__e15e9162b38246c56c8d7ca12c83c648(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract the table content from a specific table in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'table_index' (which table to check)

    Returns:
        Dict: Contains 'col_count', 'row_count', 'headers', and 'rows' data.
              Returns empty dict if table doesn't exist.
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        table_index = config.get('table_index', 0)
        if len(doc.tables) <= table_index:
            return {}
        table = doc.tables[table_index]
        col_count = len(table.columns)
        row_count = len(table.rows)
        rows = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            rows.append(row_data)
        headers = rows[0] if rows else []
        return {'col_count': col_count, 'row_count': row_count, 'headers': headers, 'rows': rows}
    except Exception as e:
        return {}
    finally:
        os.unlink(tmp_path)

def get_chrome_search_engine__26a050291c58f47430d6ebca143c45ff(env, config):
    """Get the default search engine name from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Search engine name (e.g., "Bing", "Google", "DuckDuckGo")
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        if not content:
            logger.error('Failed to read Chrome Preferences file')
            return 'Google'
        data = json.loads(content)
        search_engine = data.get('default_search_provider_data', {}).get('template_url_data', {}).get('short_name', 'Google')
        logger.info(f'Default search engine: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error reading search engine: {e}')
        return 'Google'

def get_gdrive_file_check__661594c9(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_chrome_fixed_font_size__87f268be1c9c07cc266e3e54a412bd8c(env, config: Dict[str, str]):
    """Get Chrome's fixed-width/monospace font size setting from Preferences file."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13})
        return webprefs
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'default_fixed_font_size': 13}

def get_docx_table_content__f06c6ff6344501bcedf4d007c342b6fd(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract table content from a DOCX file with section context.

    This getter reads a DOCX file and extracts all tables along with their section context
    (the heading/section they appear under).

    Args:
        env: DesktopEnv instance with controller
        config: Configuration dict with 'path' key pointing to the DOCX file on VM

    Returns:
        Dictionary with tables data and section information, or None if error occurs
        Format: {
            'num_tables': int,
            'tables': [
                {
                    'num_rows': int,
                    'num_cols': int,
                    'data': [[cell_text, ...], ...],
                    'section': str,  # The section/heading this table appears under
                    'preceding_heading': str  # The immediate heading before the table
                },
                ...
            ]
        }
    """
    try:
        file_path = config.get('path', '')
        if not file_path:
            logger.error('No path specified in config')
            return None
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {file_path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            tables_data = []
            current_section = ''
            last_heading = ''
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    para = Paragraph(element, doc)
                    text = para.text.strip()
                    if para.style and para.style.name and ('Heading' in para.style.name):
                        last_heading = text
                        if 'Heading 1' in para.style.name or text.lower() in ['main results', 'introduction', 'methods', 'results', 'conclusion']:
                            current_section = text
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    table_info = {'num_rows': len(table.rows), 'num_cols': len(table.columns) if table.rows else 0, 'data': [], 'section': current_section, 'preceding_heading': last_heading}
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_info['data'].append(row_data)
                    tables_data.append(table_info)
            result = {'num_tables': len(tables_data), 'tables': tables_data}
            logger.info(f'Extracted {len(tables_data)} tables with section context')
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting DOCX table content: {e}')
        return None

def get_csv_table_mech__9629fd7e(env, config):
    """Get CSV table content for mech task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/MECH-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_docx_table_first_row__f2b7a56f41e0323869ae9541f5efcf40(env, config: Dict[str, Any]) -> List[str]:
    """Extract the first row (header) from a specific table in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'table_index'

    Returns:
        List[str]: List of cell texts from the first row, or empty list if table doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        table_index = config.get('table_index', 0)
        if len(doc.tables) <= table_index:
            return []
        table = doc.tables[table_index]
        if len(table.rows) == 0:
            return []
        first_row = table.rows[0]
        return [cell.text.strip() for cell in first_row.cells]
    except Exception as e:
        return []
    finally:
        os.unlink(tmp_path)

def get_docx_table_row2_formatting__6412b9b5(env, config):
    """Check strikethrough in table row 2 (index 1)."""
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'has_strike': False}
    doc = Document(io.BytesIO(file_bytes))
    if len(doc.tables) == 0:
        return {'has_strike': False}
    table = doc.tables[0]
    if len(table.rows) < 2:
        return {'has_strike': False}
    row = table.rows[1]
    has_strike = False
    for cell in row.cells:
        for para in cell.paragraphs:
            for run in para.runs:
                if run.font.strike:
                    has_strike = True
                    break
    return {'has_strike': has_strike}

def get_docx_table_data__7e712d1d6a81a74414ae2e1785ec944e(env, config):
    """Extract table data and formatting from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'rows' (list of lists), 'num_columns' (int), and 'formatting' (dict)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'rows': [], 'num_columns': 0, 'formatting': {}}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return {'rows': [], 'num_columns': 0, 'formatting': {}}
        table = doc.tables[0]
        table_data = []
        num_columns = 0
        formatting_info = {'has_bold': False, 'has_borders': False, 'num_formatted_cells': 0, 'total_cells': 0}
        for (row_idx, row) in enumerate(table.rows):
            row_data = []
            for (cell_idx, cell) in enumerate(row.cells):
                cell_text = cell.text.strip()
                formatting_info['total_cells'] += 1
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if run.bold:
                            formatting_info['has_bold'] = True
                            formatting_info['num_formatted_cells'] += 1
                            break
                if table.style is not None:
                    formatting_info['has_borders'] = True
                try:
                    if '.' in cell_text:
                        row_data.append(float(cell_text))
                    else:
                        row_data.append(int(cell_text))
                except (ValueError, AttributeError):
                    row_data.append(cell_text)
            table_data.append(row_data)
            num_columns = max(num_columns, len(row_data))
        return {'rows': table_data, 'num_columns': num_columns, 'formatting': formatting_info}
    finally:
        os.unlink(tmp_path)

def get_chrome_download_location__518a9ee4(env, config: Dict[str, Any]) -> str:
    """
    Get the default download directory from Chrome preferences.

    Args:
        env: Environment object with vm_platform and controller
        config: Configuration dict (not used in this function)

    Returns:
        str: The default_directory from download preferences, or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception(f'Unsupported operating system: {os_type}')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        default_directory = data.get('download', {}).get('default_directory', '')
        logger.info(f'Retrieved download location: {default_directory}')
        return default_directory
    except Exception as e:
        logger.error(f'Error retrieving download location: {e}')
        return ''

def get_third_party_cookies_blocked__75ed94d8e5e2ff083f81c247646b0b38(env, config: Dict[str, str]):
    """
    Check if third-party cookies are blocked in Chrome.
    Returns "true" if third-party cookies are blocked, "false" otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        cookie_controls_mode = data.get('profile', {}).get('cookie_controls_mode', None)
        if cookie_controls_mode is not None:
            return 'true' if cookie_controls_mode >= 1 else 'false'
        block_third_party = data.get('profile', {}).get('block_third_party_cookies', False)
        return 'true' if block_third_party else 'false'
    except Exception as e:
        logger.error(f'Error checking third-party cookies setting: {e}')
        return 'false'

def get_docx_table_data__9453c274a20969d63590381e37c06d85(env, config):
    """Extract table data and formatting from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'rows', 'num_rows', 'num_cols', and 'formatting'
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'rows': [], 'num_rows': 0, 'num_cols': 0, 'formatting': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return {'rows': [], 'num_rows': 0, 'num_cols': 0, 'formatting': []}
        table = doc.tables[0]
        table_data = []
        formatting_data = []
        max_cols = 0
        for row in table.rows:
            row_data = []
            row_formatting = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                try:
                    if '.' in cell_text:
                        row_data.append(float(cell_text))
                    else:
                        row_data.append(int(cell_text))
                except (ValueError, AttributeError):
                    row_data.append(cell_text)
                cell_format = {}
                try:
                    if cell._element.tcPr is not None:
                        shd = cell._element.tcPr.find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd')
                        if shd is not None and shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill'):
                            fill_color = shd.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}fill')
                            cell_format['bg_color'] = fill_color
                    if cell.paragraphs:
                        para = cell.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                            if run.font.color and run.font.color.rgb:
                                cell_format['font_color'] = str(run.font.color.rgb)
                            if run.font.bold is not None:
                                cell_format['bold'] = run.font.bold
                            if run.font.italic is not None:
                                cell_format['italic'] = run.font.italic
                            if run.font.name:
                                cell_format['font_name'] = run.font.name
                            if run.font.size:
                                cell_format['font_size'] = run.font.size.pt
                except Exception:
                    pass
                row_formatting.append(cell_format)
            table_data.append(row_data)
            formatting_data.append(row_formatting)
            max_cols = max(max_cols, len(row_data))
        return {'rows': table_data, 'num_rows': len(table_data), 'num_cols': max_cols, 'formatting': formatting_data}
    finally:
        os.unlink(tmp_path)

def get_table_count_only__db5e4b5e(env, config):
    """
    Get table count and bulleted list information from document.

    Returns:
        dict with:
        - table_count: number of tables in document
        - has_bulleted_list: whether document has bulleted lists
        - bulleted_list_items: list of bulleted items found
        - list_location_index: paragraph index where list starts
        - table_positions: list of paragraph indices where tables appear
        - tables_with_consonants: list of table indices containing the expected consonants
    """
    file_path = config.get('path')
    if not file_path:
        return {'table_count': 0, 'has_bulleted_list': False, 'bulleted_list_items': [], 'list_location_index': -1, 'table_positions': [], 'tables_with_consonants': []}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        return {'table_count': 0, 'has_bulleted_list': False, 'bulleted_list_items': [], 'list_location_index': -1, 'table_positions': [], 'tables_with_consonants': []}
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        table_count = len(doc.tables)
        expected_consonants = ['p', 'b', 't', 'd', 'k', 'c', 'g']
        table_positions = []
        tables_with_consonants = []
        for (idx, element) in enumerate(doc.element.body):
            if element.tag.endswith('tbl'):
                table_positions.append(idx)
        for (table_idx, table) in enumerate(doc.tables):
            table_text = ''
            for row in table.rows:
                for cell in row.cells:
                    table_text += cell.text.lower() + ' '
            consonants_in_table = set()
            for consonant in expected_consonants:
                if f' {consonant} ' in table_text or f' {consonant},' in table_text or f' {consonant}-' in table_text or table_text.startswith(f'{consonant} ') or table_text.startswith(f'{consonant},') or table_text.startswith(f'{consonant}-'):
                    consonants_in_table.add(consonant)
            if len(consonants_in_table) >= 5:
                tables_with_consonants.append(table_idx)
                logger.info(f'Table {table_idx} contains consonants: {consonants_in_table}')
        bulleted_lists = []
        list_location_indices = []
        for (para_idx, paragraph) in enumerate(doc.paragraphs):
            pPr = paragraph._element.pPr
            if pPr is not None:
                numPr = pPr.find(qn('w:numPr'))
                if numPr is not None:
                    numId = numPr.find(qn('w:numId'))
                    if numId is not None:
                        num_id_val = numId.get(qn('w:val'))
                        if num_id_val == '0':
                            continue
                        try:
                            num_part = doc.part.numbering_part
                            if num_part:
                                num_elem = num_part.element
                                is_bullet = False
                                for num in num_elem.findall(qn('w:num')):
                                    if num.get(qn('w:numId')) == num_id_val:
                                        abstractNumId = num.find(qn('w:abstractNumId'))
                                        if abstractNumId is not None:
                                            abstract_num_id = abstractNumId.get(qn('w:val'))
                                            for abstractNum in num_elem.findall(qn('w:abstractNum')):
                                                if abstractNum.get(qn('w:abstractNumId')) == abstract_num_id:
                                                    for lvl in abstractNum.findall(qn('w:lvl')):
                                                        numFmt = lvl.find(qn('w:numFmt'))
                                                        if numFmt is not None:
                                                            fmt_val = numFmt.get(qn('w:val'))
                                                            if fmt_val == 'bullet':
                                                                is_bullet = True
                                                            break
                                                    break
                                        break
                                if is_bullet:
                                    text = paragraph.text.strip()
                                    if text:
                                        bulleted_lists.append(text)
                                        if para_idx not in list_location_indices:
                                            list_location_indices.append(para_idx)
                        except Exception as e:
                            logger.debug(f'Could not check numFmt: {e}')
        has_bulleted_list = len(bulleted_lists) > 0
        list_location_index = min(list_location_indices) if list_location_indices else -1
        os.unlink(tmp_path)
        result = {'table_count': table_count, 'has_bulleted_list': has_bulleted_list, 'bulleted_list_items': bulleted_lists, 'list_location_index': list_location_index, 'table_positions': table_positions, 'tables_with_consonants': tables_with_consonants}
        logger.info(f'Found {table_count} tables at positions {table_positions}, bulleted list: {has_bulleted_list}, items: {bulleted_lists}, list location: {list_location_index}, tables with consonants: {tables_with_consonants}')
        return result
    except Exception as e:
        logger.error(f'Error: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return {'table_count': 0, 'has_bulleted_list': False, 'bulleted_list_items': [], 'list_location_index': -1, 'table_positions': [], 'tables_with_consonants': []}

def get_chrome_startup_urls__715fc21b1c707dd79fc5ab9b6e4df514(env, config: Dict[str, str]):
    """
    Get the list of startup URLs configured in Chrome.
    Returns a list of URLs that Chrome opens on startup.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        startup_urls = data.get('session', {}).get('startup_urls', [])
        logger.info(f'Current startup URLs: {startup_urls}')
        return startup_urls
    except Exception as e:
        logger.error(f'Error getting startup URLs: {e}')
        return []

def get_extension_files_count__014e651b(env, config: Dict[str, Any]):
    """Count how many required extension files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' and 'required_files' parameters

    Returns:
        int: Number of required files that exist
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    required_files = config.get('required_files', ['manifest.json', 'hello.html', 'popup.js'])
    count = 0
    for file_name in required_files:
        file_path = f'{extension_path}/{file_name}'
        try:
            result = env.controller.run_bash_script(f"test -f {file_path} && echo 'exists' || echo 'not_found'", timeout=10)
            output = result.get('output', '').strip()
            if output == 'exists':
                count += 1
        except Exception as e:
            logger.error(f'Error checking file {file_name}: {e}')
    return count

def get_chrome_font_size__d996a712(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 6.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_gdrive_pdf_file__9e5c8fdc(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_webext_dir__2ec2a9a8(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_extension_manifest_version__b7035c23581ef82d8725cee1b3aa987f(env, config: dict):
    """
    Get the manifest_version of an installed unpacked extension by its path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        int: Extension manifest_version or 0 if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        target_path = config.get('extension_path', '')
        for (ext_id, ext_data) in all_extensions.items():
            if ext_data.get('path') == target_path:
                manifest = ext_data.get('manifest', {})
                return manifest.get('manifest_version', 0)
        return 0
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return 0

def get_webext_dir__60d86cd5(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_chrome_tabs_and_bookmarks__6d016e82(env, config: dict):
    """Get Chrome tab count and bookmark information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: {
            "open_tab_count": int,
            "bookmarks": list
        }
    """
    result = {'open_tab_count': 0, 'bookmarks': []}
    try:
        tabs = get_open_tabs_info(env, {})
        if tabs:
            result['open_tab_count'] = len(tabs)
            logger.info(f"Found {result['open_tab_count']} open tabs")
    except Exception as e:
        logger.error(f'Error getting tabs: {e}')
    try:
        bookmarks = get_bookmarks(env, {})
        if bookmarks:
            result['bookmarks'] = bookmarks
            logger.info(f'Found {len(bookmarks)} bookmarks')
    except Exception as e:
        logger.error(f'Error getting bookmarks: {e}')
    return result

def get_gdrive_file_list__fd7bb036(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_vim_tabstop_config__76faceaeb09651393a582783b81e5798(env, config):
    """Check if tabstop is set to 4 in Vim configuration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with command parameters

    Returns:
        str: Command output indicating if tabstop is set
    """
    command = config.get('command', 'bash eval.sh')
    shell = config.get('shell', True)
    vm_ip = env.vm_ip
    port = env.server_port
    import requests
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': shell})
    if response.status_code == 200:
        return response.json()['output']
    else:
        logger.error('Failed to get vim config. Status code: %d', response.status_code)
        return None

def get_webext_manifest__70dff70b667529be05282f3babd2fe6e(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract manifest.json from a web extension project directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to manifest.json on VM)

    Returns:
        Dict containing the manifest JSON, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest file from VM: {path}')
            return None
        manifest = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded manifest from {path}')
        return manifest
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading manifest from {path}: {e}')
        return None

def get_csv_table_math__bdcc0312(env, config):
    """Get CSV table content for math task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/MATH-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_block_third_party_cookies__980fd38b(env, config: Dict[str, str]):
    """Check if third-party cookies are blocked in Chrome preferences."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        block_third_party_cookies = profile.get('block_third_party_cookies', False)
        return 'true' if block_third_party_cookies else 'false'
    except Exception as e:
        logger.error(f'Error: {e}')
        return 'false'

def get_url_text__5525d1c8(env, config):
    """Extract URL from docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text from the document
    """
    path = config.get('path', '/home/user/Desktop/signup_url.docx')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return ''
    import os
    cache_path = os.path.join(env.cache_dir, 'signup_url.docx')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        text = ' '.join([p.text.strip() for p in doc.paragraphs if p.text.strip()])
        return text
    except:
        return ''

def get_webext_dir__c1aa21b2(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_macys_url_parse__7dbecc9b465f893baa757d13deeefa17(env, config: Dict[str, str]):
    """
    Parse Macy's product URL for men's medium short-sleeve shirts with 50% discount.
    Variation 0: men's medium-size short-sleeve shirts with 50% discount.
    """
    result = {}
    active_tab_url = get_active_url_from_accessTree(env, config)
    if active_tab_url is None:
        return None
    parsed = urlparse(active_tab_url)
    path = unquote(parsed.path)
    result['mens_clothing'] = True if 'mens-clothing' in path else None
    path_parts = path.strip('/').split('/')
    key_value_json = {}
    shirts_flag = False
    short_sleeve_flag = False
    if 'shirts' in path:
        shirts_flag = True
    if 'short-sleeve' in path:
        short_sleeve_flag = True
    for i in range(len(path_parts) - 1):
        if ',' in path_parts[i] and ',' in path_parts[i + 1]:
            keys = [k.strip() for k in path_parts[i].split(',')]
            values = [v.strip() for v in path_parts[i + 1].split(',')]
            for (k, v) in zip(keys, values):
                if k == 'Price_discount_range':
                    key_value_json[k] = [item.strip() for item in v.split('|')] if v else None
                else:
                    key_value_json[k] = v if v else None
                if k == 'Product_department' and (v == 'shirts' or v == 'Shirts' or v == 'Shirt'):
                    shirts_flag = True
                if k == 'Sleeve_length' and (v == 'short-sleeve' or v == 'Short Sleeve'):
                    short_sleeve_flag = True
            break
    for field in ['Men_regular_size_t', 'Price_discount_range']:
        if field not in key_value_json:
            key_value_json[field] = None
    result['shirts'] = shirts_flag if shirts_flag else None
    result['short_sleeve'] = short_sleeve_flag if short_sleeve_flag else None
    for key in config.get('parse_keys', []):
        if key in key_value_json:
            if key == 'Price_discount_range':
                if key_value_json[key] is not None and '50_PERCENT_ off & more' in key_value_json[key] and (not '30_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                    result[key] = '50_PERCENT_ off & more'
                else:
                    result[key] = 'not_50_PERCENT_ off & more'
            else:
                result[key] = key_value_json[key]
    return result

def get_extension_icon_exists__db044e4c(env, config: Dict[str, Any]):
    """Check if the extension icon file exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' and 'icon_name' parameters

    Returns:
        bool: True if icon file exists, False otherwise
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    icon_name = config.get('icon_name', 'hello_extensions.png')
    icon_path = f'{extension_path}/{icon_name}'
    try:
        result = env.controller.run_bash_script(f"test -f {icon_path} && echo 'exists' || echo 'not_found'", timeout=10)
        output = result.get('output', '').strip()
        return output == 'exists'
    except Exception as e:
        logger.error(f'Error checking extension icon: {e}')
        return False

def get_gdrive_file_check__42b1e128(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_csv_table_cs4y__5116177c(env, config):
    """Get CSV table content for cs4y task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/CS-p4y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_docx_table_row_count__ff8ca82a4f791a2f5346fef371999b2c(env, config: Dict[str, Any]) -> int:
    """Extract the row count from a specific table in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'table_index' (which table to check)

    Returns:
        int: Number of rows in the specified table, or 0 if table doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return 0
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        table_index = config.get('table_index', 0)
        if len(doc.tables) <= table_index:
            return 0
        table = doc.tables[table_index]
        return len(table.rows)
    except Exception as e:
        return 0
    finally:
        os.unlink(tmp_path)

def get_csv_file_and_chrome_tab__4e1aca7b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if CSV file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    csv_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.csv'
    file_exists = False
    try:
        result = env.controller.get_file(csv_path)
        file_exists = result is not None and len(result) > 0
        logger.info(f'CSV file check: path={csv_path}, exists={file_exists}')
    except Exception as e:
        logger.warning(f'Error checking CSV file: {e}')
        file_exists = False
    chrome_tabs = []
    try:
        from desktop_env.evaluators.getters.chrome import get_open_tabs_info
        tabs_info = get_open_tabs_info(env, {})
        if tabs_info:
            chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
        logger.info(f'Chrome tabs: {chrome_tabs}')
    except Exception as e:
        logger.warning(f'Error getting Chrome tabs: {e}')
        chrome_tabs = []
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': csv_path}

def get_googledrive_file_list__dd26081a(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a Google Drive folder

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_query: Query string to find the folder

    Returns:
        List of filenames in the folder, or empty list if folder not found
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            logger.warning(f'Folder not found with query: {folder_query}')
            return []
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        filenames = [f['title'] for f in file_list]
        logger.info(f'Found {len(filenames)} files in folder: {filenames}')
        return sorted(filenames)
    except Exception as e:
        logger.error(f'Error getting Google Drive file list: {e}')
        return []

def get_extension_manifest__1edda3bd4444a8fb1554a227ff178017(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract manifest.json content from browser extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to manifest.json

    Returns:
        Dictionary containing manifest.json content, or empty dict if file not found
    """
    path = config.get('path', '')
    if not path.startswith('/home/user/Projects/'):
        logger.warning(f'Manifest path is not in /home/user/Projects: {path}')
        return {}
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest.json from VM: {path}')
            return {}
        manifest_data = json.loads(file_bytes.decode('utf-8'))
        if manifest_data:
            logger.info(f'Successfully parsed manifest.json from {path}')
        return manifest_data
    except json.JSONDecodeError as e:
        logger.error(f'JSON decode error in manifest.json: {e}')
        return {}
    except Exception as e:
        logger.error(f'Error reading manifest.json from {path}: {e}')
        return {}

def get_default_web_browser__865c6c015a0aa9b8277b205e904948d3(env, config: dict):
    """Gets the default web browser application.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: The default web browser .desktop file name (e.g., 'firefox.desktop')
    """
    import requests
    os_type = env.vm_platform
    if os_type == 'Linux':
        mime_types = ['text/html', 'text/xml', 'application/xhtml+xml', 'application/xml', 'x-scheme-handler/http', 'x-scheme-handler/https']
        apps = []
        vm_ip = env.vm_ip
        port = env.server_port
        for mime_type in mime_types:
            command = ['xdg-mime', 'query', 'default', mime_type]
            response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
            if response.status_code == 200:
                app = response.json().get('output', '').strip()
                if app:
                    apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    else:
        raise Exception('Unsupported operating system', os_type)

def get_webext_manifest__d1a4f844666cd6560872c801fdefe60a(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract manifest.json from a web extension project directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to manifest.json on VM)

    Returns:
        Dict containing the manifest JSON, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest file from VM: {path}')
            return None
        manifest = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded manifest from {path}')
        return manifest
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading manifest from {path}: {e}')
        return None

def get_table_position__f4eb9543(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_extension_description__2224641eff529090c9cc8ad62127e176(env, config: dict):
    """
    Get the description of an installed unpacked extension by its path.

    This getter reads Chrome's Preferences file to find the extension at the
    specified path and extracts the description from its manifest. It includes
    retry logic to handle timing issues where Chrome may not have updated the
    Preferences file immediately after extension installation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension description or empty string if not found/not enabled
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    target_path = config.get('extension_path', '')
    max_retries = 3
    retry_delay = 2
    for attempt in range(max_retries):
        try:
            content = env.controller.get_file(preference_file_path)
            data = json.loads(content)
            all_extensions = data.get('extensions', {}).get('settings', {})
            for (ext_id, ext_data) in all_extensions.items():
                if ext_data.get('path') == target_path:
                    state = ext_data.get('state', 0)
                    if state != 1:
                        logger.warning(f'Extension found at {target_path} but not enabled (state={state})')
                        return ''
                    manifest = ext_data.get('manifest', {})
                    description = manifest.get('description', '')
                    if description:
                        return description
                    else:
                        logger.warning(f'Extension found at {target_path} but description is empty')
                        return ''
            if attempt < max_retries - 1:
                logger.info(f'Extension not found at {target_path}, retrying ({attempt + 1}/{max_retries})...')
                time.sleep(retry_delay)
            else:
                logger.warning(f'Extension not found at {target_path} after {max_retries} attempts')
                return ''
        except Exception as e:
            logger.error(f'Error reading Chrome Preferences file: {e}')
            if attempt < max_retries - 1:
                time.sleep(retry_delay)
            else:
                return ''
    return ''

def get_extension_version__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, str]:
    """Get version information for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict mapping extension names to their version strings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extensions_version = {}
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            name = ext_data.get('manifest', {}).get('name', '')
            version = ext_data.get('manifest', {}).get('version', '')
            if name:
                extensions_version[name] = version
        return extensions_version
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file for extension versions: {e}')
        return {}

def get_csv_table_bio__a4ea4d07(env, config):
    """Get CSV table content for biological sciences task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/BIO-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_chrome_policy_page_info__f3b19d1e(env, config: Dict[str, str]) -> Optional[Dict]:
    """
    Get active tab URL and visible text content to validate policy page.
    Uses get_active_url_from_accessTree to avoid page reload,
    then separately gets the page content and extracts visible text.

    Args:
        env: Environment object
        config: Configuration dict (optional parameters for URL extraction)

    Returns:
        Dict with 'url' and 'content' keys (content is visible text only), or None if failed
    """
    try:
        from desktop_env.evaluators.getters.chrome import get_active_url_from_accessTree
        active_url = get_active_url_from_accessTree(env, config)
        if not active_url:
            logger.error('[POLICY_PAGE_INFO] Failed to get active tab URL')
            return None
        logger.info(f'[POLICY_PAGE_INFO] Active tab URL: {active_url}')
        host = env.vm_ip
        port = env.chromium_port
        remote_debugging_url = f'http://{host}:{port}'
        from playwright.sync_api import sync_playwright
        import time
        max_retries = 2
        timeout_ms = 30000
        for attempt in range(max_retries):
            try:
                logger.info(f'[POLICY_PAGE_INFO] Attempt {attempt + 1}/{max_retries} to get page content')
                with sync_playwright() as p:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    contexts = browser.contexts
                    if not contexts:
                        logger.error('[POLICY_PAGE_INFO] No browser contexts found')
                        return None
                    pages = contexts[0].pages
                    if not pages:
                        logger.error('[POLICY_PAGE_INFO] No pages found in context')
                        return None
                    active_page = None
                    for page in pages:
                        try:
                            if page.url == active_url or active_url in page.url:
                                active_page = page
                                break
                        except Exception as e:
                            logger.warning(f'[POLICY_PAGE_INFO] Error checking page: {e}')
                            continue
                    if not active_page and pages:
                        active_page = pages[0]
                    if not active_page:
                        logger.error('[POLICY_PAGE_INFO] Could not find active page')
                        return None
                    html_content = active_page.content()
                    visible_text = extract_visible_text(html_content)
                    logger.info(f'[POLICY_PAGE_INFO] Successfully retrieved page content (HTML length: {len(html_content)}, visible text length: {len(visible_text)})')
                    browser.close()
                    return {'url': active_url, 'content': visible_text}
            except Exception as e:
                logger.error(f'[POLICY_PAGE_INFO] Attempt {attempt + 1} failed: {e}')
                if attempt < max_retries - 1:
                    logger.info('[POLICY_PAGE_INFO] Retrying in 2 seconds...')
                    time.sleep(2)
                else:
                    logger.error(f'[POLICY_PAGE_INFO] All {max_retries} attempts failed')
                    return None
        return None
    except Exception as e:
        logger.error(f'[POLICY_PAGE_INFO] Error in get_chrome_policy_page_info__f3b19d1e: {e}')
        return None

def get_docx_table_info__a4c465e5(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_gdrive_text_file__a48d3d2ba069ee77685e4821041681b9(env, config: dict):
    """Get text file content from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'settings_file' and 'query_list'

    Returns:
        str: Content of the text file from Google Drive
    """
    from desktop_env.evaluators.getters.file import get_googledrive_file
    file_bytes = get_googledrive_file(env, config)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        print(f'Error decoding file: {e}')
        return ''

def get_gdrive_pdf_file__9d82c3b2(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_chrome_setting_value__df4ecdef(env, config: Dict[str, str]):
    """
    Get the Homepage URL setting from Chrome preferences.
    Checks both startup_urls and restore_on_startup settings.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        startup_urls = []
        session_data = data.get('session', {})
        if isinstance(session_data, dict):
            startup_urls = session_data.get('startup_urls', [])
        restore_on_startup = session_data.get('restore_on_startup', 1) if isinstance(session_data, dict) else 1
        logger.info(f'[CHROME_SETTING] Retrieved startup_urls: {startup_urls}')
        logger.info(f'[CHROME_SETTING] Retrieved restore_on_startup: {restore_on_startup}')
        return {'setting_value': startup_urls, 'restore_on_startup': restore_on_startup}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': [], 'restore_on_startup': 1}

def get_chrome_min_font_size__f6a59c3b(env, config: dict):
    """Extract minimum font size from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        int: Minimum font size in pixels
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        min_font_size = data.get('webkit', {}).get('webprefs', {}).get('minimum_font_size', 0)
        logger.info(f'Chrome minimum font size: {min_font_size}')
        return int(min_font_size)
    except Exception as e:
        logger.error(f'Error getting Chrome minimum font size: {e}')
        return 0

def get_chrome_font_size__caf9b8e1(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 9.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_gdrive_pdf_file__75da4280(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_chrome_min_font_size__904067aa(env, config):
    """Get Chrome's minimum font size setting from Preferences file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Minimum font size settings including 'minimum_font_size'
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        minimum_font_size = webprefs.get('minimum_font_size', 0)
        return {'minimum_font_size': minimum_font_size}
    except Exception as e:
        logger.error(f'Error getting minimum font size: {e}')
        return {'minimum_font_size': 0}

def get_gdrive_pdf_info__bb91e7693f2a30704f1d1cc79be73950(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Get PDF file from Google Drive and extract basic information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with Google Drive settings and path

    Returns:
        dict: Dictionary with 'exists', 'page_count', 'file_size' keys, or None if file not found
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        from pypdf import PdfReader
    except ImportError as e:
        logger.error(f'Missing required library: {e}')
        return None
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    path_list = config.get('path', [])
    if not path_list:
        logger.warning('No path specified in config')
        return None
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for (idx, folder_or_file) in enumerate(path_list):
            is_folder = idx < len(path_list) - 1
            if is_folder:
                search = f"title = '{folder_or_file}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false and '{parent_id}' in parents"
            else:
                search = f"title = '{folder_or_file}' and trashed = false and '{parent_id}' in parents"
            filelist = drive.ListFile({'q': search}).GetList()
            if len(filelist) == 0:
                logger.info(f"File/folder '{folder_or_file}' not found in Google Drive")
                return {'exists': False, 'page_count': 0, 'file_size': 0}
            file = filelist[0]
            parent_id = file['id']
        cache_path = os.path.join(env.cache_dir, config.get('dest', 'temp.pdf'))
        os.makedirs(os.path.dirname(cache_path) if os.path.dirname(cache_path) else env.cache_dir, exist_ok=True)
        file.GetContentFile(cache_path, mimetype=file.get('mimeType', 'application/pdf'))
        if not os.path.exists(cache_path):
            logger.warning('Failed to download file from Google Drive')
            return {'exists': False, 'page_count': 0, 'file_size': 0}
        file_size = os.path.getsize(cache_path)
        reader = PdfReader(cache_path)
        page_count = len(reader.pages)
        logger.info(f'PDF found in Google Drive: {page_count} pages, {file_size} bytes')
        return {'exists': True, 'page_count': page_count, 'file_size': file_size}
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return {'exists': False, 'page_count': 0, 'file_size': 0}

def get_html_exists__7b7b0b2f(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Check if an HTML file exists on the VM and verify it contains valid HTML content.

    Args:
        env: Environment object with controller to access VM
        config: Configuration dict with 'path' key pointing to HTML file on VM

    Returns:
        dict: {
            'exists': bool indicating if file exists,
            'is_valid_html': bool indicating if content is valid HTML,
            'is_non_empty': bool indicating if file has content
        }
        Returns None if there's an error accessing the file
    """
    html_path_on_vm = config['path']
    try:
        file_content = env.controller.get_file(html_path_on_vm)
        if file_content is None:
            return {'exists': False, 'is_valid_html': False, 'is_non_empty': False}
        is_non_empty = len(file_content) > 0
        is_valid_html = False
        if is_non_empty:
            content_str = file_content.decode('utf-8') if isinstance(file_content, bytes) else file_content
            content_lower = content_str.strip().lower()
            html_indicators = [content_lower.startswith('<!doctype html'), content_lower.startswith('<html'), '<html>' in content_lower, '<html ' in content_lower, '<!doctype' in content_lower]
            is_valid_html = any(html_indicators)
        return {'exists': True, 'is_valid_html': is_valid_html, 'is_non_empty': is_non_empty}
    except Exception as e:
        return None

def get_extension_manifest_name__e5eabc9a(env, config: Dict[str, Any]):
    """Extract the extension name from manifest.json in the unpacked extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension name from manifest.json, or None if not found
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    manifest_path = f'{extension_path}/manifest.json'
    try:
        content = env.controller.get_file(manifest_path)
        if content:
            manifest_data = json.loads(content)
            return manifest_data.get('name', None)
        return None
    except Exception as e:
        logger.error(f'Error reading extension manifest: {e}')
        return None

def get_html_file_and_chrome_tab__3cb1a587(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if HTML file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict containing html_path from expected rules

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    html_path = config.get('expected', {}).get('rules', {}).get('html_path', '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.html')
    result = env.controller.get_file(html_path)
    file_exists = result is not None and len(result) > 0
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    tabs_info = get_open_tabs_info(env, {})
    chrome_tabs = []
    if tabs_info:
        chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': html_path}

def get_webext_dir__88e21052(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_chrome_setting_value__fca61186(env, config: Dict[str, str]):
    """
    Get the Fixed-width font size setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['webkit', 'webprefs', 'default_fixed_font_size']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, 13)
            else:
                setting_value = 13
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': 13}

def get_googledrive_eml_files(env, config: Dict[str, Any]) -> List[str]:
    """
    Get all .eml files from a specified Google Drive folder.

    Downloads all .eml files from the specified folder path and returns
    their local file paths.

    Args:
        env: Environment object (provides cache_dir for downloads)
        config: Configuration dict containing:
            - settings_file: Path to Google Drive authentication settings
            - folder_path: List representing folder path, e.g., ['emails']

    Returns:
        List[str]: List of local file paths to downloaded .eml files,
                   or empty list if folder not found or no .eml files exist
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false and '{parent_id}' in parents"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                logger.info(f"Folder '{folder_name}' not found in parent '{parent_id}'")
                return []
            folder = filelist[0]
            parent_id = folder['id']
        query = f"title contains '.eml' and trashed = false and '{parent_id}' in parents"
        filelist = drive.ListFile({'q': query}).GetList()
        downloaded_paths = []
        for (idx, file) in enumerate(filelist):
            if not file['title'].endswith('.eml'):
                continue
            dest_path = os.path.join(env.cache_dir, f"gdrive_eml_{idx}_{file['title']}")
            try:
                file.GetContentFile(dest_path, mimetype=file.get('mimeType', 'text/plain'))
                downloaded_paths.append(dest_path)
                logger.info(f"Downloaded .eml file: {file['title']} to {dest_path}")
            except Exception as e:
                logger.error(f"Failed to download {file['title']}: {e}")
        return downloaded_paths
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_vim_hlsearch_config__5438ce42ea45fa77c023ecd730e398a5(env, config):
    """Check if search highlighting is enabled in Vim configuration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with command parameters

    Returns:
        str: Command output indicating if hlsearch is set
    """
    command = config.get('command', 'bash eval.sh')
    shell = config.get('shell', True)
    vm_ip = env.vm_ip
    port = env.server_port
    import requests
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': shell})
    if response.status_code == 200:
        return response.json()['output']
    else:
        logger.error('Failed to get vim config. Status code: %d', response.status_code)
        return None

def get_docx_last_table_dims__d15d9273(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_vim_tabstop_check__e2e649c5246b836f874ad28b723333bc(env, config: dict):
    """
    Check if .vimrc file contains 'set tabstop=4' and 'set expandtab' configurations.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        String output from the check script
    """
    script = '\nif [ -f ~/.vimrc ]; then\n    # Check for uncommented \'set tabstop=4\' line\n    if grep -E \'^[[:space:]]*set tabstop=4\' ~/.vimrc > /dev/null; then\n        tabstop_found=true\n    else\n        tabstop_found=false\n    fi\n\n    # Check for uncommented \'set expandtab\' line\n    if grep -E \'^[[:space:]]*set expandtab\' ~/.vimrc > /dev/null; then\n        expandtab_found=true\n    else\n        expandtab_found=false\n    fi\n\n    if [ "$tabstop_found" = true ] && [ "$expandtab_found" = true ]; then\n        echo "The File Has Set Tabstop=4 And Expandtab!"\n    elif [ "$tabstop_found" = true ]; then\n        echo "The File Has Set Tabstop=4 But Missing Expandtab!"\n    else\n        echo "The File Does Not Have Set Tabstop=4!"\n    fi\nelse\n    echo "The .vimrc File Does Not Exist!"\nfi\n'
    result = env.controller.run_bash_script(script, timeout=10)
    output = result.get('output', '').strip()
    logger.info(f'Vim tabstop check result: {output}')
    return output

def get_html_result__b5ca523a54cc6ae837c08b60601febd6(env, config: dict):
    """
    Extract both command history and HTML file count from archive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' and 'path' keys

    Returns:
        dict: {
            'command_output': str,
            'html_count': int,
            'archive_exists': bool
        }
    """
    command_result = env.controller.run_bash_script(config['command'], timeout=30)
    command_output = command_result.get('output', '') if command_result else ''
    html_count = 0
    archive_exists = False
    count_command = ['/bin/bash', '-c', f"if [ -f {config['path']} ]; then tar -tzf {config['path']} 2>/dev/null | grep '\\.html$' | wc -l; else echo '0'; fi"]
    count_result = env.controller.run_bash_script(count_command, timeout=30)
    if count_result:
        count_output = count_result.get('output', '0').strip()
        try:
            html_count = int(count_output)
            if html_count > 0:
                archive_exists = True
        except ValueError:
            html_count = 0
    return {'command_output': command_output, 'html_count': html_count, 'archive_exists': archive_exists}

def get_chrome_notifications_blocked__4f8fdc8b(env, config: dict):
    """
    Check if Chrome's notifications are blocked for all sites.
    This checks the 'profile.default_content_setting_values.notifications' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if notifications are blocked (value=2), "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        notifications_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('notifications', 0)
        return 'true' if notifications_setting == 2 else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_webext_dir__bba937aa(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_extension_name__a4b0d616259699774cd5ee2679c8a96c(env, config: dict):
    """
    Get the name of an installed unpacked extension by its path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension name or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        target_path = config.get('extension_path', '')
        for (ext_id, ext_data) in all_extensions.items():
            if ext_data.get('path') == target_path:
                manifest = ext_data.get('manifest', {})
                return manifest.get('name', '')
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return ''

def get_chrome_default_encoding__e3f03b16(env, config: Dict[str, str]):
    """Get Chrome default text encoding setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with default_encoding value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        encoding = data.get('intl', {}).get('charset_default', 'ISO-8859-1')
        return {'default_encoding': encoding}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'default_encoding': 'ISO-8859-1'}

def get_gdrive_twitter_receipt__12ff505f1cc77038855fb34142173dbc(env, config: Dict[str, Any]) -> Optional[str]:
    """Get Twitter receipt email file from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_name: Name of the folder containing file (default: 'emails')
            - filename: Exact filename to search for
            - dest: Local destination filename

    Returns:
        Local filepath where file was downloaded, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    filename = config.get('filename')
    dest = config.get('dest', 'pred.eml')
    if not filename:
        logger.error('No filename specified in config')
        return None
    auth = GoogleAuth(settings_file=settings_file)
    drive = GoogleDrive(auth)
    folder_query = f"title = '{folder_name}' and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
    folder_list: GoogleDriveFileList = drive.ListFile({'q': folder_query}).GetList()
    if len(folder_list) == 0:
        logger.warning(f"Folder '{folder_name}' not found in Google Drive")
        return None
    folder_id = folder_list[0]['id']
    file_query = f"title = '{filename}' and '{folder_id}' in parents"
    file_list: GoogleDriveFileList = drive.ListFile({'q': file_query}).GetList()
    if len(file_list) == 0:
        logger.warning(f"File '{filename}' not found in folder '{folder_name}'")
        return None
    file: GoogleDriveFile = file_list[0]
    try:
        file.GetContentFile(dest, mimetype=file['mimeType'])
        return dest
    except Exception as e:
        logger.error(f"Failed to download '{filename}': {e}")
        return None

def get_unpacked_extension_count__0f84311e(env, config):
    """
    Get the count of unpacked extensions loaded in Chrome.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        int: Number of unpacked extensions
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        unpacked_count = 0
        all_extensions = data.get('extensions', {}).get('settings', {})
        for extension_id in all_extensions.keys():
            extension_info = all_extensions[extension_id]
            if 'path' in extension_info:
                unpacked_count += 1
                logger.info(f"Found unpacked extension: {extension_id} at path: {extension_info['path']}")
        logger.info(f'Total unpacked extensions found: {unpacked_count}')
        return unpacked_count
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return 0

def get_docx_table_info__2341f5ea(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_chrome_setting_value__0717aaed(env, config: Dict[str, str]):
    """
    Get the Password saving setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['profile', 'password_manager_enabled']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, None)
            else:
                setting_value = None
                break
        if setting_value is None:
            logger.warning(f'[CHROME_SETTING] Setting not found, defaulting to True')
            setting_value = True
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': True}

def get_chrome_setting_value__fd2ee811(env, config: Dict[str, str]):
    """
    Get the Autofill setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data.get('autofill', {}).get('enabled', False)
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': False}

def get_chrome_camera_blocked__d2e0663c(env, config: dict):
    """
    Check if Chrome's camera access is blocked for all sites.
    This checks the 'profile.default_content_setting_values.media_stream_camera' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if camera is blocked (value=2), "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        camera_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('media_stream_camera', 0)
        return 'true' if camera_setting == 2 else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_tabs_and_bookmarks__7a5a7856f1b642a4ade91ca81ca0f263000420251221151547(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get both open tabs and bookmarks data.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with 'tabs' and 'bookmarks' keys
    """
    logger.info('Getting tabs and bookmarks data')
    tabs = get_open_tabs_info(env, {})
    logger.info(f'Retrieved {(len(tabs) if tabs else 0)} open tabs')
    bookmarks = get_bookmarks(env, {})
    logger.info(f'Retrieved bookmarks data')
    return {'tabs': tabs or [], 'bookmarks': bookmarks or {}}

def get_docx_table_size__e7b89121046aaa7adf88d3e2fb20c6ab(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get size (rows and columns) from tables in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        Dict with table information: {'tables': [{'rows': int, 'columns': int}, ...]}
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return {'tables': []}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'tables': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        tables_info = []
        for table in doc.tables:
            tables_info.append({'rows': len(table.rows), 'columns': len(table.columns)})
        return {'tables': tables_info}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_googledrive_docx__a60c9fbb05dbe8be7b4812ad58480d11(env, config: Dict[str, Any]) -> Any:
    """Get DOCX file from Google Drive forms/ folder.

    This getter reuses the existing get_googledrive_file function.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with Google Drive query parameters

    Returns:
        Path to downloaded file, or None if not found
    """
    from desktop_env.evaluators.getters.chrome import get_googledrive_file
    return get_googledrive_file(env, config)

def get_csv_table_phys__5e1848a9(env, config):
    """Get CSV table content for phys task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/PHYS-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_chrome_fixed_font_size__b7f59f23(env, config: Dict[str, str]):
    """Get Chrome default fixed (monospace) font size setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with default_fixed_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13})
        return {'default_fixed_font_size': webprefs.get('default_fixed_font_size', 13)}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'default_fixed_font_size': 13}

def get_html_file_info__48e4da460d6f3e132e4d3cc48ac9ca24(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if an HTML file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' key indicating if file exists
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'HTML file not found at {path}')
        return {'exists': False, 'path': path}
    logger.info(f'HTML file found at {path} ({len(file_bytes)} bytes)')
    return {'exists': True, 'path': path, 'size': len(file_bytes)}

def get_extension_source_type__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, str]:
    """Get source type (webstore/unpacked) for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict mapping extension names to their source type ('webstore' or 'unpacked')
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extensions_source = {}
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            name = ext_data.get('manifest', {}).get('name', '')
            from_webstore = ext_data.get('from_webstore', False)
            source_type = 'webstore' if from_webstore else 'unpacked'
            if name:
                extensions_source[name] = source_type
        return extensions_source
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file for extension source: {e}')
        return {}

def get_extension_version__f107c64c583fbfb012ea91827c8f61e3(env, config: dict):
    """
    Get the version of an installed unpacked extension by its path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension version or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        target_path = config.get('extension_path', '')
        for (ext_id, ext_data) in all_extensions.items():
            if ext_data.get('path') == target_path:
                manifest = ext_data.get('manifest', {})
                return manifest.get('version', '')
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return ''

def get_chrome_search_engine_changed__6c30ec39(env, config: dict):
    """
    Check if Chrome's default search engine has been changed from Google.
    This checks the 'default_search_provider_data' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if search engine is not Google, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_provider = data.get('default_search_provider_data', {}).get('template_url_data', {})
        short_name = search_provider.get('short_name', 'Google')
        return 'true' if short_name != 'Google' else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_gdrive_file_list__5876322dac1e30402171f6bcd2edb019(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a Google Drive folder by path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_path: List of folder names forming the path (e.g., ['emails'])

    Returns:
        List of filenames found in the specified folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"'{parent_id}' in parents and title='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
            file_list = drive.ListFile({'q': query}).GetList()
            if not file_list:
                logger.warning(f"Folder '{folder_name}' not found in path {folder_path}")
                return []
            parent_id = file_list[0]['id']
        query = f"'{parent_id}' in parents and trashed=false"
        file_list = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in file_list]
        logger.info(f'Found {len(filenames)} files in Google Drive folder: {filenames}')
        return sorted(filenames)
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_googledrive_file_list__2b29a90e(env, config: Dict[str, Any]) -> List[Dict[str, str]]:
    """Get list of files with metadata.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: List of file info dicts
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                return []
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        files = [{'title': f['title'], 'mimeType': f.get('mimeType', '')} for f in filelist if f['mimeType'] != 'application/vnd.google-apps.folder']
        return files
    except Exception as e:
        logger.error(f'Error: {e}')
        return []

def get_gdrive_file_check__adf6c9b9(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_chrome_min_logical_font__b612e026(env, config: Dict[str, str]):
    """Get Chrome minimum logical font size setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with minimum_logical_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {'minimum_logical_font_size': 6})
        return {'minimum_logical_font_size': webprefs.get('minimum_logical_font_size', 6)}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'minimum_logical_font_size': 6}

def get_chrome_safe_browsing__2b56847e(env, config):
    """
    Get the Safe Browsing status from Chrome preferences.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        bool: True if Safe Browsing is enabled (either standard or enhanced), False otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        safebrowsing = data.get('safebrowsing', {})
        is_enhanced = bool(safebrowsing.get('enhanced', False))
        is_enabled = bool(safebrowsing.get('enabled', False))
        result = is_enhanced or is_enabled
        logger.info(f'Safe Browsing status - Enhanced: {is_enhanced}, Enabled: {is_enabled}, Result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error: {e}')
        return False

def get_chrome_default_search__8a4fe4a4(env, config: Dict[str, Any]) -> str:
    """
    Get the default search engine template URL from Chrome preferences.

    Args:
        env: Environment object with vm_platform and controller
        config: Configuration dict (not used in this function)

    Returns:
        str: The template_url from default_search_provider_data, or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception(f'Unsupported operating system: {os_type}')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        template_url = data.get('default_search_provider_data', {}).get('template_url', '')
        logger.info(f'Retrieved template_url: {template_url}')
        return template_url
    except Exception as e:
        logger.error(f'Error retrieving default search engine template URL: {e}')
        return ''

def get_extension_icon__0c7a03ab(env, config: Dict[str, Any]):
    """
    Get the default_icon value from a Chrome extension's manifest in Preferences file.

    Args:
        env: Environment object
        config: Configuration dict with 'extension_name' key

    Returns:
        str: The default_icon filename from the extension's manifest, or empty string if not found
    """
    extension_name = config.get('extension_name', '')
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for (extension_id, extension_data) in all_extensions.items():
            manifest = extension_data.get('manifest', {})
            name = manifest.get('name', '')
            if name == extension_name:
                action = manifest.get('action', {})
                default_icon = action.get('default_icon', '')
                if not default_icon:
                    browser_action = manifest.get('browser_action', {})
                    default_icon = browser_action.get('default_icon', '')
                logger.info(f"Found extension '{extension_name}' with default_icon: {default_icon}")
                return default_icon
        logger.warning(f"Extension '{extension_name}' not found in Chrome Preferences")
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file for extension icon: {e}')
        return ''

def get_extension_popup__8aee113e(env, config: Dict[str, Any]) -> str:
    """
    Get the default_popup configuration from a Chrome extension's manifest.json.

    This function retrieves the extension's manifest and extracts the default_popup
    value from the action configuration.

    Args:
        env: Environment object
        config: Configuration dict containing:
            - extension_name (str): The name of the extension to check

    Returns:
        str: The default_popup HTML file path (e.g., "hello.html"), or None if not found
    """
    logger.info(f"[EXTENSION_POPUP] Starting to get extension popup for extension: {config.get('extension_name', 'N/A')}")
    extension_name = config.get('extension_name', '')
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        logger.error(f'[EXTENSION_POPUP] Unsupported operating system: {os_type}')
        return None
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        logger.info(f'[EXTENSION_POPUP] Successfully loaded Chrome Preferences')
        all_extensions = data.get('extensions', {}).get('settings', {})
        logger.info(f'[EXTENSION_POPUP] Found {len(all_extensions)} installed extensions')
        for (extension_id, extension_data) in all_extensions.items():
            try:
                manifest = extension_data.get('manifest', {})
                manifest_name = manifest.get('name', '')
                logger.debug(f'[EXTENSION_POPUP] Checking extension ID {extension_id}: {manifest_name}')
                if manifest_name == extension_name:
                    logger.info(f'[EXTENSION_POPUP] Found matching extension: {extension_name}')
                    state = extension_data.get('state', 0)
                    if state != 1:
                        logger.warning(f'[EXTENSION_POPUP] Extension found but is not enabled (state: {state})')
                        return None
                    action = manifest.get('action', {})
                    browser_action = manifest.get('browser_action', {})
                    default_popup = action.get('default_popup') or browser_action.get('default_popup')
                    if default_popup:
                        logger.info(f'[EXTENSION_POPUP] Successfully extracted default_popup: {default_popup}')
                        return default_popup
                    else:
                        logger.warning(f'[EXTENSION_POPUP] Extension found but no default_popup configured')
                        return None
            except Exception as e:
                logger.debug(f'[EXTENSION_POPUP] Error processing extension {extension_id}: {e}')
                continue
        logger.warning(f"[EXTENSION_POPUP] Extension '{extension_name}' not found in installed extensions")
        logger.info(f"[EXTENSION_POPUP] Available extensions: {[data.get('manifest', {}).get('name', 'Unknown') for data in all_extensions.values()]}")
        return None
    except Exception as e:
        logger.error(f'[EXTENSION_POPUP] Error reading Chrome Preferences file: {e}')
        return None

def get_chrome_setting_value__5936c775(env, config: Dict[str, str]):
    """
    Get the Minimum font size setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['webkit', 'webprefs', 'minimum_font_size']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, 0)
            else:
                setting_value = 0
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': 0}

def get_docx_table_structure__8586f38ff35ab9e474999a7f667bbc31(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get table structure (number of tables, rows, and columns) from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        Dict with table information: {'tables': [{'rows': int, 'columns': int}, ...]}
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return {'tables': []}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'tables': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        tables_info = []
        for table in doc.tables:
            tables_info.append({'rows': len(table.rows), 'columns': len(table.columns)})
        return {'tables': tables_info}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_html_file_and_chrome_tab__d6487d33(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if HTML file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with rules containing html_path

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    html_path = config.get('rules', {}).get('html_path', '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.html')
    result = env.controller.get_file(html_path)
    file_exists = False
    if result is not None and len(result) > 0:
        content_lower = result.lower()
        file_exists = '<html' in content_lower or '<!doctype html' in content_lower
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    tabs_info = get_open_tabs_info(env, {})
    chrome_tabs = []
    if tabs_info:
        chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': html_path}

def get_chrome_experiments_exact_match__d7481e87a8a6dde50669ce517c215edf(env, config: Dict[str, str]):
    """
    Get enabled Chrome experiments and return them as a list.
    This getter is used to check exact matches of enabled experiments.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, but required by framework)

    Returns:
        List of enabled experiment names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Local State'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enabled_labs_experiments = data.get('browser', {}).get('enabled_labs_experiments', [])
        experiment_names = [exp.split('@')[0] for exp in enabled_labs_experiments]
        return experiment_names
    except Exception as e:
        logger.error(f'Error getting enabled experiments: {e}')
        return []

def get_block_third_party_cookies__62856a905c85744fe1f63d2f847937c7(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to block third-party cookies.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: "true" if third-party cookies are blocked, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile_data = data.get('profile', {})
        block_third_party = profile_data.get('block_third_party_cookies', False)
        cookie_controls = profile_data.get('cookie_controls_mode', 0)
        is_blocked = block_third_party or cookie_controls == 2
        return 'true' if is_blocked else 'false'
    except Exception as e:
        logger.error(f'Error checking third-party cookie blocking: {e}')
        return 'false'

def get_chrome_clear_on_exit__da95206e(env, config: dict):
    """
    Check if Chrome's 'Clear cookies and site data when you quit Chrome' setting is enabled.
    This checks the 'profile.exit_type' and related preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if clear on exit is enabled, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        exit_type = data.get('profile', {}).get('exit_type', 0)
        clear_lso = data.get('browser', {}).get('clear_lso_data_enabled', False)
        return 'true' if exit_type == 1 or clear_lso else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_extension_description__221bd7ef80e48e21d62e7d38635cf26c(env, config):
    """Get description of a specific extension by name.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_name' key

    Returns:
        str: Description of the extension, or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    target_name = config.get('extension_name', '')
    if not target_name:
        return ''
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            if 'manifest' in ext_data:
                name = ext_data['manifest'].get('name', '')
                if name.lower() == target_name.lower():
                    description = ext_data['manifest'].get('description', '')
                    return description
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return ''

def get_recreation_echocanyon_html__dbcda1af4d2b810312479cfb54a15ab6(env, config: Dict[str, Any]):
    """
    Get HTML content from recreation.gov page for Echo Canyon search.
    This is a custom getter for the Echo Canyon task variation.
    Verifies that Echo Canyon location is displayed and availability table exists.
    """
    logger.info(f'[RECREATION_ECHOCANYON] Starting recreation.gov page processing for Echo Canyon')
    logger.debug(f'[RECREATION_ECHOCANYON] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 3
    timeout_ms = 60000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_ECHOCANYON] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_ECHOCANYON] Successfully connected to existing Chrome instance')
                except Exception as e:
                    logger.warning(f'[RECREATION_ECHOCANYON] Failed to connect to existing Chrome instance: {e}')
                    logger.info(f'[RECREATION_ECHOCANYON] Starting new Chrome instance...')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337', '--no-sandbox']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    logger.info(f"[RECREATION_ECHOCANYON] Starting browser with command: {' '.join(command)}")
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup' + '/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_ECHOCANYON] Successfully connected to new Chrome instance')
                if len(browser.contexts) == 0 or len(browser.contexts[0].pages) == 0:
                    logger.error(f'[RECREATION_ECHOCANYON] No active pages found')
                    return None
                page = browser.contexts[0].pages[0]
                current_url = page.url
                logger.info(f'[RECREATION_ECHOCANYON] Current URL: {current_url}')
                content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                selector = config.get('selector', 'class')
                class_name = config.get('class', 'camp-sortable-column-header')
                order = config.get('order', '2')
                result = {'has_table_header': False, 'has_echo_canyon_text': False, 'echo_canyon_in_url': False, 'page_title': '', 'url': current_url}
                if selector == 'class':
                    elements = soup.find_all(class_=class_name)
                    logger.info(f"[RECREATION_ECHOCANYON] Found {len(elements)} elements with class '{class_name}'")
                    if len(elements) > int(order):
                        target_element = elements[int(order)]
                        result['has_table_header'] = True
                        logger.info(f'[RECREATION_ECHOCANYON] Successfully found table header element at position {order}')
                    else:
                        logger.warning(f'[RECREATION_ECHOCANYON] Not enough elements found (need {int(order) + 1}, got {len(elements)})')
                page_text = soup.get_text().lower()
                if 'echo canyon' in page_text:
                    result['has_echo_canyon_text'] = True
                    logger.info(f"[RECREATION_ECHOCANYON] Found 'Echo Canyon' text in page content")
                else:
                    logger.warning(f"[RECREATION_ECHOCANYON] 'Echo Canyon' text not found in page content")
                url_lower = current_url.lower()
                if 'echo' in url_lower and 'canyon' in url_lower:
                    result['echo_canyon_in_url'] = True
                    logger.info(f"[RECREATION_ECHOCANYON] Found 'Echo Canyon' reference in URL")
                else:
                    logger.info(f"[RECREATION_ECHOCANYON] No 'Echo Canyon' reference in URL")
                title_tag = soup.find('title')
                if title_tag:
                    result['page_title'] = title_tag.get_text().strip()
                    logger.info(f"[RECREATION_ECHOCANYON] Page title: {result['page_title']}")
                for tag in ['h1', 'h2', 'h3', 'h4']:
                    headings = soup.find_all(tag)
                    for h in headings:
                        if 'echo canyon' in h.get_text().lower():
                            logger.info(f"[RECREATION_ECHOCANYON] Found 'Echo Canyon' in {tag} heading")
                            result['has_echo_canyon_text'] = True
                            break
                logger.info(f'[RECREATION_ECHOCANYON] Final result: {result}')
                return result
        except Exception as e:
            logger.error(f'[RECREATION_ECHOCANYON] Attempt {attempt + 1} failed: {e}')
            if attempt < max_retries - 1:
                logger.info(f'[RECREATION_ECHOCANYON] Retrying in 2 seconds...')
                time.sleep(2)
            else:
                logger.error(f'[RECREATION_ECHOCANYON] All retries exhausted')
                return None
    return None

def get_chrome_disable_safe_browsing__19371712(env, config: Dict[str, str]):
    """
    Check if Chrome Safe Browsing is completely disabled (both enhanced and standard).
    Returns 'true' if disabled, 'false' if any safe browsing is enabled.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        safebrowsing = data.get('safebrowsing', {})
        is_enhanced = bool(safebrowsing.get('enhanced', False))
        is_enabled = bool(safebrowsing.get('enabled', False))
        return 'true' if not is_enhanced and (not is_enabled) else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome Safe Browsing status: {e}')
        return 'false'

def get_docx_table_info__7c0843e4(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_docx_table_content__121c468593688d91f702e8d5fdd1e9f9(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract table content from a DOCX file and Excel source data.

    This getter reads a DOCX file and extracts all tables, along with the source
    Excel file data for verification.

    Args:
        env: DesktopEnv instance with controller
        config: Configuration dict with 'path' key pointing to the DOCX file on VM

    Returns:
        Dictionary with tables data and Excel source data, or None if error occurs
        Format: {
            'num_tables': int,
            'tables': [
                {
                    'num_rows': int,
                    'num_cols': int,
                    'data': [[cell_text, ...], ...],
                    'context_before': str  # Text before table for section detection
                },
                ...
            ],
            'excel_data': {
                'gpt-4': {'chrome': val, 'vscode': val, 'avg': val}
            }
        }
    """
    try:
        file_path = config.get('path', '')
        if not file_path:
            logger.error('No path specified in config')
            return None
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {file_path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            tables_data = []
            for table in doc.tables:
                table_info = {'num_rows': len(table.rows), 'num_cols': len(table.columns) if table.rows else 0, 'data': [], 'context_before': ''}
                context_paragraphs = []
                found_table = False
                for element in doc.element.body:
                    if element.tag.endswith('tbl'):
                        if element == table._element:
                            found_table = True
                            break
                    elif element.tag.endswith('p'):
                        context_paragraphs.append(element)
                        if len(context_paragraphs) > 5:
                            context_paragraphs.pop(0)
                context_text = []
                for para_elem in context_paragraphs:
                    para_text = ''.join((node.text for node in para_elem.iter() if hasattr(node, 'text') and node.text))
                    if para_text.strip():
                        context_text.append(para_text.strip())
                table_info['context_before'] = '\n'.join(context_text[-3:])
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_info['data'].append(row_data)
                tables_data.append(table_info)
            excel_data = {}
            excel_path = '/home/user/Documents/awesome-desktop/expe-results.xlsx'
            try:
                excel_bytes = env.controller.get_file(excel_path)
                if excel_bytes:
                    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_excel:
                        tmp_excel.write(excel_bytes)
                        tmp_excel_path = tmp_excel.name
                    try:
                        wb = openpyxl.load_workbook(tmp_excel_path, data_only=True)
                        ws = wb.active
                        headers = []
                        gpt4_row = None
                        for (row_idx, row) in enumerate(ws.iter_rows(values_only=True), start=1):
                            if not any(row):
                                continue
                            first_cell = str(row[0]).lower() if row[0] else ''
                            if 'model' in first_cell or row_idx == 1:
                                headers = [str(cell).lower().strip() if cell else '' for cell in row]
                            if 'gpt-4' in first_cell or 'gpt4' in first_cell.replace('-', ''):
                                gpt4_row = row
                                break
                        if gpt4_row and headers:
                            gpt4_data = {}
                            for (idx, header) in enumerate(headers):
                                if idx < len(gpt4_row):
                                    value = gpt4_row[idx]
                                    if header in ['chrome', 'vscode', 'avg', 'average']:
                                        try:
                                            gpt4_data[header] = float(value) if value is not None else None
                                        except (ValueError, TypeError):
                                            gpt4_data[header] = str(value) if value is not None else None
                                    else:
                                        gpt4_data[header] = str(value) if value is not None else None
                            excel_data['gpt-4'] = gpt4_data
                            logger.info(f'Extracted Excel data for GPT-4: {gpt4_data}')
                    finally:
                        if os.path.exists(tmp_excel_path):
                            os.unlink(tmp_excel_path)
                else:
                    logger.warning('Could not load Excel file for verification')
            except Exception as e:
                logger.warning(f'Could not extract Excel data: {e}')
            result = {'num_tables': len(tables_data), 'tables': tables_data, 'excel_data': excel_data}
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting DOCX table content: {e}')
        return None

def get_chrome_do_not_track__a858c66a(env, config: Dict[str, str]):
    """Get Chrome Do Not Track status from preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Do Not Track settings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        dnt_enabled = data.get('enable_do_not_track', False)
        return {'enabled': dnt_enabled}
    except Exception as e:
        logger.error(f'Error getting Do Not Track status: {e}')
        return {'enabled': False}

def get_chrome_enhanced_safe_browsing__a9ae68b06063b362f29e78385fb296bc(env, config):
    """Get the enhanced safe browsing setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: "true" if enhanced safe browsing is enabled, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        if not content:
            logger.error('Failed to read Chrome Preferences file')
            return 'false'
        data = json.loads(content)
        is_enhanced = data.get('safebrowsing', {}).get('enhanced', False)
        logger.info(f'Enhanced safe browsing enabled: {is_enhanced}')
        return 'true' if is_enhanced else 'false'
    except Exception as e:
        logger.error(f'Error reading enhanced safe browsing setting: {e}')
        return 'false'

def get_chrome_location_blocked__d6e72d6c(env, config: Dict[str, str]):
    """
    Check if Chrome location access is blocked by default.
    Returns 'true' if location is blocked (denied by default), 'false' otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        default_content_settings = profile.get('default_content_setting_values', {})
        geolocation_setting = default_content_settings.get('geolocation', 3)
        is_blocked = geolocation_setting == 2
        return 'true' if is_blocked else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome location settings: {e}')
        return 'false'

def get_table_count_only__d0985f12(env, config):
    """Get table count and verify merged vowel table content from document."""
    file_path = config.get('path')
    if not file_path:
        return {'table_count': 0, 'has_merged_table': False, 'merged_table_vowels': [], 'merged_table_rows': 0, 'original_vowel_tables_found': 0, 'total_original_rows': 0, 'merged_table_position': -1, 'first_vowel_table_position': -1, 'content_verification_passed': False}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        return {'table_count': 0, 'has_merged_table': False, 'merged_table_vowels': [], 'merged_table_rows': 0, 'original_vowel_tables_found': 0, 'total_original_rows': 0, 'merged_table_position': -1, 'first_vowel_table_position': -1, 'content_verification_passed': False}
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        table_count = len(doc.tables)
        vowels = ['a', 'e', 'i', 'o', 'u']
        first_four_vowel_tables = []
        all_tables_info = []
        for (idx, table) in enumerate(doc.tables):
            if len(table.rows) == 0:
                all_tables_info.append({'index': idx, 'type': 'empty', 'rows': 0})
                continue
            table_text = ''
            header_text = ''
            if table.rows:
                for cell in table.rows[0].cells:
                    header_text += ' ' + cell.text.lower().strip()
            for row in table.rows:
                for cell in row.cells:
                    table_text += ' ' + cell.text.lower().strip()
            single_vowel_match = None
            for vowel in vowels:
                header_and_start = (header_text + ' ' + table_text[:150]).lower()
                pattern1 = re.search('\\b' + vowel + '\\s+grapheme|\\bgrapheme\\s+' + vowel + '\\b', header_and_start)
                pattern2 = re.search('sound[:\\s]+' + vowel + '\\b|' + vowel + '\\s+sound', header_and_start)
                pattern3 = re.search('^\\s*' + vowel + '\\s*$|^\\s*' + vowel.upper() + '\\s*$', header_text)
                if pattern1 or pattern2 or pattern3:
                    other_vowels = [v for v in vowels if v != vowel]
                    other_vowel_patterns = 0
                    for other_v in other_vowels:
                        if re.search('\\b' + other_v + '\\s+grapheme|\\bgrapheme\\s+' + other_v + '\\b', header_and_start):
                            other_vowel_patterns += 1
                    if other_vowel_patterns == 0:
                        single_vowel_match = vowel
                        break
            if single_vowel_match:
                table_info = {'index': idx, 'type': 'single_vowel', 'vowel': single_vowel_match, 'rows': len(table.rows), 'header': header_text, 'content_sample': table_text[:300]}
                all_tables_info.append(table_info)
                if len(first_four_vowel_tables) < 4:
                    first_four_vowel_tables.append(table_info)
                    logger.info(f"Found vowel table #{len(first_four_vowel_tables)} for '{single_vowel_match}' at position {idx} with {len(table.rows)} rows")
            else:
                vowel_patterns_count = 0
                vowels_found = []
                for vowel in vowels:
                    if re.search('\\b' + vowel + '\\s+grapheme|\\bgrapheme\\s+' + vowel + '\\b', table_text):
                        vowel_patterns_count += 1
                        vowels_found.append(vowel)
                table_info = {'index': idx, 'type': 'multi_vowel' if vowel_patterns_count >= 4 else 'other', 'rows': len(table.rows), 'vowels_found': vowels_found, 'vowel_patterns_count': vowel_patterns_count, 'content': table_text}
                all_tables_info.append(table_info)
        first_vowel_table_position = first_four_vowel_tables[0]['index'] if first_four_vowel_tables else -1
        expected_total_rows = sum((info['rows'] for info in first_four_vowel_tables))
        merged_table_candidate = None
        for table_info in all_tables_info:
            if table_info['type'] == 'multi_vowel':
                position_diff = abs(table_info['index'] - first_vowel_table_position) if first_vowel_table_position >= 0 else 999
                if table_info['vowel_patterns_count'] >= 4 and position_diff <= 3:
                    content_matches = 0
                    for orig_table in first_four_vowel_tables:
                        orig_content = orig_table['content_sample']
                        words = [w for w in orig_content.split() if len(w) > 2 and w not in ['the', 'and', 'for', 'grapheme', 'sound', 'pattern']]
                        if words:
                            sample_word = words[0] if len(words) > 0 else ''
                            if sample_word and sample_word in table_info['content']:
                                content_matches += 1
                    if merged_table_candidate is None or content_matches > merged_table_candidate.get('content_matches', 0):
                        merged_table_candidate = {'index': table_info['index'], 'rows': table_info['rows'], 'vowels': table_info['vowels_found'], 'content': table_info['content'], 'content_matches': content_matches, 'position_diff': position_diff}
                        logger.info(f"Found merged table candidate at position {table_info['index']} (distance {position_diff} from first vowel table), {content_matches} content matches")
        remaining_vowel_tables = []
        for table_info in all_tables_info:
            if table_info['type'] == 'single_vowel':
                is_first_four = any((t['index'] == table_info['index'] for t in first_four_vowel_tables))
                if is_first_four:
                    remaining_vowel_tables.append(table_info)
        has_merged_table = False
        merged_table_vowels = []
        merged_table_rows = 0
        merged_table_position = -1
        content_verification_passed = False
        if merged_table_candidate:
            has_merged_table = True
            merged_table_vowels = sorted(merged_table_candidate['vowels'])
            merged_table_rows = merged_table_candidate['rows']
            merged_table_position = merged_table_candidate['index']
            content_verification_passed = merged_table_candidate['content_matches'] >= 3
            logger.info(f"Merged table verification: position={merged_table_position}, rows={merged_table_rows}, vowels={merged_table_vowels}, content_matches={merged_table_candidate['content_matches']}")
        os.unlink(tmp_path)
        result = {'table_count': table_count, 'has_merged_table': has_merged_table, 'merged_table_vowels': merged_table_vowels, 'merged_table_rows': merged_table_rows, 'original_vowel_tables_found': len(remaining_vowel_tables), 'total_original_rows': expected_total_rows, 'merged_table_position': merged_table_position, 'first_vowel_table_position': first_vowel_table_position, 'content_verification_passed': content_verification_passed}
        logger.info(f'Result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return {'table_count': 0, 'has_merged_table': False, 'merged_table_vowels': [], 'merged_table_rows': 0, 'original_vowel_tables_found': 0, 'total_original_rows': 0, 'merged_table_position': -1, 'first_vowel_table_position': -1, 'content_verification_passed': False}

def get_chrome_min_font_enforce__0d902d89d190f0f84e2f19e726b6d1ea(env, config: Dict[str, str]):
    """
    Get the minimum font size setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Dictionary containing minimum_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        minimum_font_size = webprefs.get('minimum_font_size', 0)
        return {'minimum_font_size': minimum_font_size}
    except Exception as e:
        logger.error(f'Error getting minimum font size: {e}')
        return {'minimum_font_size': 0}

def get_gdrive_file_metadata__b2463eb9(env, config: dict):
    """Download a file from Google Drive and read its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'settings_file', 'query', and 'dest'

    Returns:
        str: File content or empty string if not found
    """
    from desktop_env.evaluators.getters.chrome import get_googledrive_file
    import os
    try:
        local_filepath = get_googledrive_file(env, config)
        if not local_filepath or not os.path.exists(local_filepath):
            return ''
        with open(local_filepath, 'r', encoding='utf-8') as f:
            content = f.read()
        return content
    except Exception:
        return ''

def get_extension_count__0fb1bf36ddef8ab5a554de56ff7f0c3d(env, config):
    """Get the count of unpacked extensions installed in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        int: Number of unpacked extensions installed
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        unpacked_count = 0
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            if ext_data.get('location') in [4, 5]:
                unpacked_count += 1
        return unpacked_count
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return 0

def get_googledrive_file_count__17298c22(env, config: Dict[str, Any]):
    """Count files in a specific Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - settings_file: path to Google Drive settings
            - folder_path: list representing path to folder (e.g., ["emails"])

    Returns:
        int: Number of files in the folder
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                logger.info(f"Folder '{folder_name}' not found")
                return 0
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        file_count = sum((1 for f in filelist if f['mimeType'] != 'application/vnd.google-apps.folder'))
        logger.info(f'Found {file_count} files in folder')
        return file_count
    except Exception as e:
        logger.error(f'Error counting Google Drive files: {e}')
        return 0

def get_docx_table_count__6ea8989b(env, config):
    """
    Get the number of tables in a DOCX file.

    Args:
        env: Environment object with controller to access VM files
        config: Configuration dict with:
            - path: Path to the .docx file on VM

    Returns:
        dict: Information about tables in the document:
            - table_count: int - Number of tables in the document
    """
    file_path = config.get('path')
    if not file_path:
        return {'table_count': -1, 'error': 'No file path provided in config'}
    try:
        local_path = env.controller.get_vm_file(file_path)
    except Exception as e:
        return {'table_count': -1, 'error': f'Failed to get file from VM: {str(e)}'}
    if not os.path.exists(local_path):
        return {'table_count': -1, 'error': f'File not found: {local_path}'}
    try:
        doc = Document(local_path)
        table_count = len(doc.tables)
        return {'table_count': table_count}
    except Exception as e:
        return {'table_count': -1, 'error': f'Error reading document: {str(e)}'}

def get_show_full_urls__96430b37(env, config: Dict[str, str]):
    """Check if Chrome is set to show full URLs in the address bar."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        show_full_urls = data.get('omnibox', {}).get('prevent_url_elisions', False)
        return 'true' if show_full_urls else 'false'
    except Exception as e:
        logger.error(f'Error: {e}')
        return 'false'

def get_table_position__199b082c(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_recreation_url_check__6dc3893f96ccd943c500af1756962de6(env, config: Dict[str, Any]):
    """
    Get the current URL to verify navigation to recreation.gov.
    """
    logger.info(f'[RECREATION_URL] Starting recreation.gov URL check')
    logger.debug(f'[RECREATION_URL] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 2
    timeout_ms = 30000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_URL] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_URL] Connected to Chrome')
                except Exception as e:
                    logger.warning(f'[RECREATION_URL] Failed to connect: {e}')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                page = browser.contexts[0].pages[0] if browser.contexts[0].pages else browser.contexts[0].new_page()
                current_url = page.url
                logger.info(f'[RECREATION_URL] Current URL: {current_url}')
                page.wait_for_load_state('networkidle', timeout=timeout_ms)
                return {'url': current_url}
        except Exception as e:
            logger.error(f'[RECREATION_URL] Attempt {attempt + 1} failed: {e}')
            if attempt == max_retries - 1:
                logger.error(f'[RECREATION_URL] All attempts failed')
                return {}
            time.sleep(2)
    return {}

def get_docx_table_row_count__7d2236cd23a11ec2c058cf74d17b7e88(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get row count information from tables in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        Dict with table information: {'tables': [{'rows': int, 'columns': int}, ...]}
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return {'tables': []}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'tables': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        tables_info = []
        for table in doc.tables:
            tables_info.append({'rows': len(table.rows), 'columns': len(table.columns)})
        return {'tables': tables_info}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_table_position__6a4e1dd4(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_googledrive_has_files__95db6624(env, config: Dict[str, Any]) -> bool:
    """Check if Google Drive folder has any files"""
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            return False
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        has_files = len(file_list) > 0
        logger.info(f'Folder has files: {has_files}')
        return has_files
    except Exception as e:
        logger.error(f'Error: {e}')
        return False

def get_extension_manifest__3d9fc0c33aef9d5f768043f561973d63(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract manifest.json content from browser extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to manifest.json

    Returns:
        Dictionary containing manifest.json content, or empty dict if file not found
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest.json from VM: {path}')
            return {}
        manifest_data = json.loads(file_bytes.decode('utf-8'))
        return manifest_data
    except json.JSONDecodeError as e:
        logger.error(f'JSON decode error in manifest.json: {e}')
        return {}
    except Exception as e:
        logger.error(f'Error reading manifest.json from {path}: {e}')
        return {}

def get_chrome_block_third_party_cookies__327d732b(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to block third-party cookies.
    Returns 'true' if third-party cookies are blocked, 'false' otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        block_third_party = profile.get('block_third_party_cookies', False)
        cookie_controls_mode = profile.get('cookie_controls_mode', 0)
        return 'true' if block_third_party or cookie_controls_mode >= 1 else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome third-party cookie blocking: {e}')
        return 'false'

def get_chrome_do_not_track__a5d5a0e3(env, config: dict):
    """
    Check if Chrome's 'Do Not Track' setting is enabled.
    This checks the 'enable_do_not_track' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if Do Not Track is enabled, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        do_not_track = data.get('enable_do_not_track', False)
        return 'true' if do_not_track else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_default_web_browser__b2d61e52(env, config: dict):
    """Gets the default web browser application.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The default web browser registered for x-scheme-handler/http
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'x-scheme-handler/http']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_backup_extension_files__8ece34a1(env, config):
    """Get list of files with .backup extension in directory.

    Args:
        env: Desktop environment
        config: Dict with 'directory' key

    Returns:
        list: Sorted list of filenames with .backup extension
    """
    directory = config.get('directory', '/home/user/Pictures')
    result = env.controller.run_bash_script(f"find {directory} -maxdepth 1 -type f -name '*.backup' | xargs -I{{}} basename {{}}", timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        return []
    files = result['output'].strip().split('\n')
    return sorted([f.strip() for f in files if f.strip()])

def get_recreation_bearlake_html__6adae5a4637b133a46055994fbaa8dd4(env, config: Dict[str, Any]):
    """
    Get HTML content from recreation.gov page for Bear Lake search.
    This getter extracts:
    - Search location verification (Bear Lake in page)
    - Reservation results presence
    - Availability information presence
    - Reservation dates (for additional verification)
    """
    logger.info(f'[RECREATION_BEARLAKE] Starting recreation.gov page processing for Bear Lake')
    logger.debug(f'[RECREATION_BEARLAKE] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 3
    timeout_ms = 60000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_BEARLAKE] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_BEARLAKE] Successfully connected to existing Chrome instance')
                except Exception as e:
                    logger.warning(f'[RECREATION_BEARLAKE] Failed to connect to existing Chrome instance: {e}')
                    logger.info(f'[RECREATION_BEARLAKE] Starting new Chrome instance...')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337', '--no-sandbox']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    logger.info(f"[RECREATION_BEARLAKE] Starting browser with command: {' '.join(command)}")
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup' + '/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_BEARLAKE] Successfully connected to new Chrome instance')
                if len(browser.contexts) == 0 or len(browser.contexts[0].pages) == 0:
                    logger.error(f'[RECREATION_BEARLAKE] No active pages found')
                    return None
                page = browser.contexts[0].pages[0]
                current_url = page.url
                logger.info(f'[RECREATION_BEARLAKE] Current URL: {current_url}')
                content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                result = {'bear_lake_found': False, 'has_results': False, 'has_availability': False, 'reservation_dates': [], 'url_contains_search': False}
                if 'bear' in current_url.lower() or 'search' in current_url.lower():
                    result['url_contains_search'] = True
                    logger.info(f'[RECREATION_BEARLAKE] URL contains search parameters')
                page_text = soup.get_text().lower()
                if 'bear lake' in page_text:
                    result['bear_lake_found'] = True
                    logger.info(f"[RECREATION_BEARLAKE] Found 'Bear Lake' in page content")
                search_inputs = soup.find_all(['input', 'textarea'])
                for input_elem in search_inputs:
                    input_value = input_elem.get('value', '').lower()
                    if 'bear lake' in input_value:
                        result['bear_lake_found'] = True
                        logger.info(f"[RECREATION_BEARLAKE] Found 'Bear Lake' in search input")
                        break
                result_indicators = [soup.find_all(class_='camp-sortable-column-header'), soup.find_all(class_='rec-grid-item'), soup.find_all(class_='campground-result'), soup.find_all('table', class_=lambda x: x and 'result' in x.lower() if x else False)]
                for indicator in result_indicators:
                    if indicator and len(indicator) > 0:
                        result['has_results'] = True
                        logger.info(f'[RECREATION_BEARLAKE] Found {len(indicator)} result elements')
                        break
                date_elements = soup.find_all(['time', 'span', 'div'], class_=lambda x: x and ('date' in x.lower() or 'available' in x.lower()) if x else False)
                for elem in date_elements:
                    date_text = elem.get_text().strip()
                    if date_text:
                        result['reservation_dates'].append(date_text)
                availability_keywords = ['available', 'open', 'vacancy', 'book']
                for keyword in availability_keywords:
                    if keyword in page_text:
                        result['has_availability'] = True
                        logger.info(f"[RECREATION_BEARLAKE] Found availability indicator: '{keyword}'")
                        break
                if result['has_results']:
                    tables = soup.find_all('table')
                    for table in tables:
                        rows = table.find_all('tr')
                        if len(rows) > 1:
                            logger.info(f'[RECREATION_BEARLAKE] Found table with {len(rows)} rows')
                logger.info(f'[RECREATION_BEARLAKE] Extraction complete: {result}')
                return result
        except Exception as e:
            logger.error(f'[RECREATION_BEARLAKE] Attempt {attempt + 1} failed: {e}')
            if attempt < max_retries - 1:
                logger.info(f'[RECREATION_BEARLAKE] Retrying in 2 seconds...')
                time.sleep(2)
            else:
                logger.error(f'[RECREATION_BEARLAKE] All retries exhausted')
                return None
    return None

def get_docx_table_structure__98ae47d7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table structure information from a docx file.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key

    Returns:
        Dict containing:
            - table_count: Number of tables in the document
            - rows: Number of rows in the first table (or 0 if no tables)
            - cols: Number of columns in the first table (or 0 if no tables)
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.error(f'Failed to get file from VM: {path}')
        return {'table_count': 0, 'rows': 0, 'cols': 0}
    cache_path = os.path.join(env.cache_dir, os.path.basename(path))
    os.makedirs(env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        table_count = len(doc.tables)
        if table_count == 0:
            return {'table_count': 0, 'rows': 0, 'cols': 0}
        first_table = doc.tables[0]
        rows = len(first_table.rows)
        cols = len(first_table.columns) if rows > 0 else 0
        logger.info(f'Document has {table_count} table(s), first table: {rows} rows x {cols} cols')
        return {'table_count': table_count, 'rows': rows, 'cols': cols}
    except Exception as e:
        logger.error(f'Error parsing docx file {path}: {e}')
        return {'table_count': 0, 'rows': 0, 'cols': 0}

def get_docx_last_table_dims__356ec724(env, config: Dict[str, Any]):
    """Get table count and dimensions of the last table in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'table_count', 'last_table_rows', 'last_table_cols' keys, or None if error
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        table_count = len(doc.tables)
        if table_count == 0:
            return {'table_count': 0, 'last_table_rows': 0, 'last_table_cols': 0}
        last_table = doc.tables[-1]
        return {'table_count': table_count, 'last_table_rows': len(last_table.rows), 'last_table_cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_chrome_extension_manifest__3d480472b35ce7003ca943ba6b2307fa(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract Chrome extension manifest.json and check if required files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - manifest_path: Path to manifest.json on VM
            - required_files: List of file paths that should exist

    Returns:
        Dict with:
            - manifest: Parsed manifest.json content (dict)
            - files_exist: Dict mapping file paths to boolean (whether they exist)
            - all_files_exist: Boolean (True if all required files exist)
    """
    manifest_path = config.get('manifest_path', '')
    required_files = config.get('required_files', [])
    result = {'manifest': {}, 'files_exist': {}, 'all_files_exist': False}
    manifest_bytes = env.controller.get_file(manifest_path)
    if not manifest_bytes:
        return result
    try:
        manifest_content = manifest_bytes.decode('utf-8')
        result['manifest'] = json.loads(manifest_content)
    except (UnicodeDecodeError, json.JSONDecodeError):
        return result
    all_exist = True
    for file_path in required_files:
        file_bytes = env.controller.get_file(file_path)
        exists = file_bytes is not None and len(file_bytes) > 0
        result['files_exist'][file_path] = exists
        if not exists:
            all_exist = False
    result['all_files_exist'] = all_exist
    return result

def get_csv_table_civil__2c0c636e(env, config):
    """Get CSV table content for civil task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/CIVIL-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_docx_table_content__a3ef22f0b6dddefc025c6922cfbdc9a0(env, config: Dict[str, Any]) -> List[List[str]]:
    """Extract table content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        List of table rows, where each row is a list of cell values
    """
    file_path = config.get('path')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_tables.append(table_data)
        return all_tables
    finally:
        os.unlink(tmp_path)

def get_gdrive_file_locations__b75faf5b6765d0d1458ed6b6d219047b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get location information of files to verify they're in the correct folder and not trashed.
    Also verifies PNG format, file size, and metadata including creation timestamp.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to check (default: 'figures')
            - expected_filenames: List of expected file names

    Returns:
        Dict with:
            - files_in_folder: List of file names in the target folder
            - files_not_trashed: List of file names that are not trashed
            - correct_location_count: Number of files in correct location and not trashed
            - file_details: Dict mapping filename to details (mime_type, size, created_date, is_recent)
            - valid_png_count: Number of files that are valid PNG images
            - recent_files_count: Number of files created recently (within last hour)
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'figures')
    expected_filenames = config.get('expected_filenames', ['1.png', '2.png', '3.png'])
    result = {'files_in_folder': [], 'files_not_trashed': [], 'correct_location_count': 0, 'file_details': {}, 'valid_png_count': 0, 'recent_files_count': 0}
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and trashed = false and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if not folder_list:
            return result
        folder_id = folder_list[0]['id']
        file_query = f"'{folder_id}' in parents"
        file_list = drive.ListFile({'q': file_query}).GetList()
        current_time = datetime.utcnow()
        one_hour_ago = current_time - timedelta(hours=1)
        for f in file_list:
            filename = f['title']
            is_trashed = f.get('labels', {}).get('trashed', False)
            mime_type = f.get('mimeType', '')
            file_size = int(f.get('fileSize', 0))
            created_date = f.get('createdDate', '')
            if not is_trashed:
                result['files_in_folder'].append(filename)
                result['files_not_trashed'].append(filename)
                is_recent = False
                try:
                    if created_date:
                        date_str = created_date.replace('Z', '').split('.')[0]
                        file_created_time = datetime.fromisoformat(date_str)
                        is_recent = file_created_time >= one_hour_ago
                except Exception:
                    is_recent = False
                result['file_details'][filename] = {'mime_type': mime_type, 'size': file_size, 'created_date': created_date, 'is_valid_png': False, 'is_recent': is_recent}
                is_valid_png = mime_type == 'image/png' and file_size > 100
                if is_valid_png:
                    result['file_details'][filename]['is_valid_png'] = True
                    result['valid_png_count'] += 1
                if is_recent:
                    result['recent_files_count'] += 1
                if filename in expected_filenames and is_valid_png and is_recent:
                    result['correct_location_count'] += 1
        return result
    except Exception as e:
        import logging
        logger = logging.getLogger('desktopenv.getter.googledrive')
        logger.error(f'Error getting Google Drive file locations: {e}')
        return result

def get_docx_table_content__fefdbfb10ebeb3b7ce54da62c0ba8cd2(env, config: Dict[str, Any]) -> List[List[str]]:
    """Extract table content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        List of table rows, where each row is a list of cell values
    """
    file_path = config.get('path')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_tables.append(table_data)
        return all_tables
    finally:
        os.unlink(tmp_path)

def get_docx_last_table_dims__b98be273(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_chrome_enhanced_safe_browsing__a6a55567(env, config: Dict[str, str]):
    """Get Chrome Enhanced Safe Browsing status from preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Enhanced safe browsing settings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enhanced_enabled = data.get('safebrowsing', {}).get('enhanced', False)
        return {'enabled': enhanced_enabled}
    except Exception as e:
        logger.error(f'Error getting enhanced safe browsing status: {e}')
        return {'enabled': False}

def get_docx_table_header__13eed80defa5785247e65a748bfcefa5(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the header row and data rows of a specific table in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'table_index' keys

    Returns:
        Dict with 'headers' (list of header texts) and 'rows' (list of row data lists)
    """
    try:
        vm_path = config['path']
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {vm_path}')
            return {'headers': [], 'rows': []}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            table_index = config.get('table_index', 0)
            if table_index >= len(doc.tables):
                logger.error(f'Table index {table_index} out of range (total tables: {len(doc.tables)})')
                return {'headers': [], 'rows': []}
            table = doc.tables[table_index]
            if len(table.rows) == 0:
                return {'headers': [], 'rows': []}
            header_row = table.rows[0]
            header_texts = [cell.text.strip() for cell in header_row.cells]
            data_rows = []
            for row in table.rows[1:]:
                row_texts = [cell.text.strip() for cell in row.cells]
                data_rows.append(row_texts)
            return {'headers': header_texts, 'rows': data_rows, 'row_count': len(data_rows)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting table data: {e}')
        return {'headers': [], 'rows': []}

def get_chrome_extension_manifest__1e834a0e2fee4e317d31e5d8fca95c5c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract Chrome extension manifest.json and check if required files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - manifest_path: Path to manifest.json on VM
            - required_files: List of file paths that should exist

    Returns:
        Dict with:
            - manifest: Parsed manifest.json content (dict)
            - files_exist: Dict mapping file paths to boolean (whether they exist)
            - all_files_exist: Boolean (True if all required files exist)
    """
    manifest_path = config.get('manifest_path', '')
    required_files = config.get('required_files', [])
    result = {'manifest': {}, 'files_exist': {}, 'all_files_exist': False}
    manifest_bytes = env.controller.get_file(manifest_path)
    if not manifest_bytes:
        return result
    try:
        manifest_content = manifest_bytes.decode('utf-8')
        result['manifest'] = json.loads(manifest_content)
    except (UnicodeDecodeError, json.JSONDecodeError):
        return result
    all_exist = True
    for file_path in required_files:
        file_bytes = env.controller.get_file(file_path)
        exists = file_bytes is not None and len(file_bytes) > 0
        result['files_exist'][file_path] = exists
        if not exists:
            all_exist = False
    result['all_files_exist'] = all_exist
    return result

def get_table_position__bbe9b961(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_extension_manifest__842d772f81b13ef554c1bda7e1c59bb7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract manifest.json content from browser extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to manifest.json

    Returns:
        Dictionary containing manifest.json content, or empty dict if file not found
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest.json from VM: {path}')
            return {}
        manifest_data = json.loads(file_bytes.decode('utf-8'))
        return manifest_data
    except json.JSONDecodeError as e:
        logger.error(f'JSON decode error in manifest.json: {e}')
        return {}
    except Exception as e:
        logger.error(f'Error reading manifest.json from {path}: {e}')
        return {}

def get_table_position__5ff0bd58(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_googledrive_filenames__a2f23245(env, config: Dict[str, Any]) -> List[str]:
    """Get list of filenames in a Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: path to Google Drive settings
            - folder_path: list representing path to folder

    Returns:
        list: List of filenames in the folder
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                return []
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in filelist if f['mimeType'] != 'application/vnd.google-apps.folder']
        return filenames
    except Exception as e:
        logger.error(f'Error getting Google Drive filenames: {e}')
        return []

def get_csv_table_med__1bf7a84e(env, config):
    """Get CSV table content for med task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of rows, each row is a list of cell values
    """
    file_path = config.get('path', '/home/user/Desktop/MED-p5y-Sheet1.csv')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    return rows

def get_chrome_default_font__ad306c7d(env, config: Dict[str, str]):
    """Get Chrome default font family setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with standard_font_family value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        fonts = data.get('webkit', {}).get('webprefs', {}).get('fonts', {}).get('standard', {})
        standard_font = fonts.get('Zyyy', 'Times New Roman')
        return {'standard_font_family': standard_font}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'standard_font_family': 'Times New Roman'}

def get_docx_table_data_rows__969da71b9e6a789c33cf79de761b11d9(env, config: Dict[str, Any]) -> int:
    """
    Get the number of data rows in a specific table (excluding header row).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'table_index' keys

    Returns:
        Number of data rows (total rows - 1), or 0 on error
    """
    try:
        vm_path = config['path']
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {vm_path}')
            return 0
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            table_index = config.get('table_index', 0)
            if table_index >= len(doc.tables):
                logger.error(f'Table index {table_index} out of range (total tables: {len(doc.tables)})')
                return 0
            table = doc.tables[table_index]
            total_rows = len(table.rows)
            data_rows = max(0, total_rows - 1)
            return data_rows
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting table data rows: {e}')
        return 0

def get_chrome_homepage_setting__32a135b2f85c9e9cdbf4b5150b97474d(env, config: Dict[str, str]):
    """
    Get the Chrome homepage setting.
    Returns the homepage URL if set, or "none" if no homepage is configured.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        homepage = data.get('homepage', '')
        logger.info(f'Current homepage: {homepage}')
        return homepage if homepage else 'none'
    except Exception as e:
        logger.error(f'Error getting homepage: {e}')
        return 'none'

def get_gdrive_file_properties__e9151d35420d1468fe1e721181658d5f(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Get properties of files in a specific Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to check (default: 'figures')

    Returns:
        List of dicts, each containing:
            - name: str (file name)
            - size: int (file size in bytes)
            - mimetype: str (MIME type)
            - is_png: bool (whether file has .png extension)
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'figures')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and trashed = false and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if not folder_list:
            return []
        folder_id = folder_list[0]['id']
        file_query = f"trashed = false and '{folder_id}' in parents"
        file_list = drive.ListFile({'q': file_query}).GetList()
        file_properties = []
        for f in file_list:
            props = {'name': f['title'], 'size': int(f.get('fileSize', 0)), 'mimetype': f.get('mimeType', ''), 'is_png': f['title'].lower().endswith('.png')}
            file_properties.append(props)
        return file_properties
    except Exception as e:
        import logging
        logger = logging.getLogger('desktopenv.getter.googledrive')
        logger.error(f'Error getting Google Drive file properties: {e}')
        return []

def get_docx_table_data__6a2751d358b86b079b53effd72604b8f(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract table data from a Word document with position information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the VM file path

    Returns:
        Dictionary containing:
        - 'tables': List of tables, where each table is a list of rows, and each row is a list of cell values
        - 'table_positions': List of dictionaries with position info for each table
            - 'index': Table index in document
            - 'after_main_results': Boolean indicating if table appears after "Main Results" heading
        Returns None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        table_positions = []
        main_results_found = False
        main_results_position = -1
        table_element_positions = {}
        for (i, table) in enumerate(doc.tables):
            table_element_positions[id(table._element)] = i
        element_index = 0
        for element in doc.element.body:
            if element.tag.endswith('p'):
                para_text = ''
                for text_element in element.itertext():
                    para_text += text_element
                para_text = para_text.strip()
                if 'Main Results' in para_text:
                    main_results_found = True
                    main_results_position = element_index
            element_index += 1
        for (table_idx, table) in enumerate(doc.tables):
            table_rows = []
            for row in table.rows:
                cell_values = [cell.text.strip() for cell in row.cells]
                table_rows.append(cell_values)
            tables_data.append(table_rows)
            table_position_in_doc = -1
            elem_idx = 0
            for element in doc.element.body:
                if element.tag.endswith('tbl'):
                    if id(element) == id(table._element):
                        table_position_in_doc = elem_idx
                        break
                elem_idx += 1
            after_main_results = main_results_found and table_position_in_doc > main_results_position
            table_positions.append({'index': table_idx, 'after_main_results': after_main_results, 'position_in_doc': table_position_in_doc})
        return {'tables': tables_data, 'table_positions': table_positions, 'main_results_found': main_results_found, 'main_results_position': main_results_position}
    except Exception as e:
        print(f'Error reading docx: {e}')
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_table_position__60d42be3(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_writer_table_font__e0ebf50c(env, config: Dict[str, Any]):
    """Extract the default table font for LibreOffice Writer.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to registrymodifications.xcu file

    Returns:
        str: The default table font name or None if not found
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    config_file_path = get_vm_file(env, config)
    if not config_file_path:
        logger.error('Failed to get config file')
        return None
    try:
        tree = ET.parse(config_file_path)
        root = tree.getroot()
        namespace = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', namespace):
            for prop in elem.findall('.//prop[@oor:name="Label"]', namespace):
                for value in prop.findall('value', namespace):
                    return value.text
        logger.warning('Default table/label font not found in config')
        return None
    except Exception as e:
        logger.error(f'Error parsing config file: {e}')
        return None

def get_extension_dir_exists__79fc6dbb(env, config: Dict[str, Any]):
    """Check if the extension directory exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        bool: True if directory exists, False otherwise
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    try:
        result = env.controller.run_bash_script(f"test -d {extension_path} && echo 'exists' || echo 'not_found'", timeout=10)
        output = result.get('output', '').strip()
        return output == 'exists'
    except Exception as e:
        logger.error(f'Error checking extension directory: {e}')
        return False

def get_gdrive_files_with_pattern__6538c56492a239416edab32324471fe6(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get files from Google Drive matching a naming pattern and extract sequential numbers.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - settings_file: Path to Google Drive settings file
            - query: Base query string
            - pattern: Regex pattern to match file names

    Returns:
        Dict with keys:
            - matching_count: int - number of files matching pattern
            - matching_files: list - names of matching files
            - file_numbers: list - extracted numbers from filenames (e.g., [1, 2, 3])
            - all_files: list - all file names found
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
    except ImportError:
        return {'matching_count': 0, 'matching_files': [], 'file_numbers': [], 'all_files': []}
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    query = config.get('query', '')
    pattern_str = config.get('pattern', '.*')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        filelist = drive.ListFile({'q': query}).GetList()
        all_files = [f['title'] for f in filelist]
        pattern = re.compile(pattern_str)
        matching_files = [name for name in all_files if pattern.match(name)]
        file_numbers = []
        number_pattern = re.compile('paper_(\\d+)\\.pdf')
        for name in matching_files:
            match = number_pattern.match(name)
            if match:
                file_numbers.append(int(match.group(1)))
        return {'matching_count': len(matching_files), 'matching_files': matching_files, 'file_numbers': sorted(file_numbers), 'all_files': all_files}
    except Exception as e:
        return {'matching_count': 0, 'matching_files': [], 'file_numbers': [], 'all_files': [], 'error': str(e)}

def get_chrome_microphone_blocked__a44d2f59(env, config: dict):
    """
    Check if Chrome's microphone access is blocked for all sites.
    This checks the 'profile.default_content_setting_values.media_stream_mic' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if microphone is blocked (value=2), "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        mic_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('media_stream_mic', 0)
        return 'true' if mic_setting == 2 else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_file_has_extension__c9bfc81c(env, config):
    """Get backup file verification info comparing original and backup files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'original_path' and 'backup_path' keys

    Returns:
        dict: Contains 'original_exists', 'backup_exists', 'original_size', 'backup_size'
    """
    original_path = config.get('original_path', '')
    backup_path = config.get('backup_path', '')
    original_bytes = env.controller.get_file(original_path)
    backup_bytes = env.controller.get_file(backup_path)
    result = {'original_exists': original_bytes is not None and len(original_bytes) > 0, 'backup_exists': backup_bytes is not None and len(backup_bytes) > 0, 'original_size': len(original_bytes) if original_bytes else 0, 'backup_size': len(backup_bytes) if backup_bytes else 0}
    return result

def get_macys_url_parse__972b80b805a61948b132613370427348(env, config: Dict[str, str]):
    """
    Parse Macy's product URL for men's small short-sleeve shirts with 20% discount.
    Variation 2: men's small short-sleeve shirts with 20% discount.
    """
    result = {}
    active_tab_url = get_active_url_from_accessTree(env, config)
    if active_tab_url is None:
        return None
    parsed = urlparse(active_tab_url)
    path = unquote(parsed.path)
    result['mens_clothing'] = True if 'mens-clothing' in path else None
    path_parts = path.strip('/').split('/')
    key_value_json = {}
    shirts_flag = False
    short_sleeve_flag = False
    long_sleeve_flag = False
    if 'shirts' in path:
        shirts_flag = True
    if 'short-sleeve' in path:
        short_sleeve_flag = True
    if 'long-sleeve' in path:
        long_sleeve_flag = True
    for i in range(len(path_parts) - 1):
        if ',' in path_parts[i] and ',' in path_parts[i + 1]:
            keys = [k.strip() for k in path_parts[i].split(',')]
            values = [v.strip() for v in path_parts[i + 1].split(',')]
            for (k, v) in zip(keys, values):
                if k == 'Price_discount_range':
                    key_value_json[k] = [item.strip() for item in v.split('|')] if v else None
                else:
                    key_value_json[k] = v if v else None
                if k == 'Product_department' and (v == 'shirts' or v == 'Shirts' or v == 'Shirt'):
                    shirts_flag = True
                if k == 'Sleeve_length':
                    if v == 'short-sleeve' or v == 'Short Sleeve':
                        short_sleeve_flag = True
                    elif v == 'long-sleeve' or v == 'Long Sleeve':
                        long_sleeve_flag = True
            break
    for field in ['Men_regular_size_t', 'Price_discount_range', 'Sleeve_length']:
        if field not in key_value_json:
            key_value_json[field] = None
    result['shirts'] = shirts_flag if shirts_flag else None
    result['short_sleeve'] = short_sleeve_flag if short_sleeve_flag else None
    result['long_sleeve'] = long_sleeve_flag if long_sleeve_flag else None
    for key in config.get('parse_keys', []):
        if key in key_value_json:
            if key == 'Price_discount_range':
                if key_value_json[key] is not None:
                    if '50_PERCENT_ off & more' in key_value_json[key]:
                        result[key] = '20_PERCENT_ off & more'
                    elif '30_PERCENT_ off & more' in key_value_json[key]:
                        result[key] = '20_PERCENT_ off & more'
                    elif '20_PERCENT_ off & more' in key_value_json[key]:
                        result[key] = '20_PERCENT_ off & more'
                    else:
                        result[key] = 'other_discount'
                else:
                    result[key] = 'no_discount'
            else:
                result[key] = key_value_json[key]
    return result

def get_docx_last_table_dims__70149acd(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_gdrive_numbered_files__aebf0a91a4be0be45ef3247f943131f2(env, config: Dict[str, Any]) -> List[Optional[str]]:
    """Get numbered email files from Google Drive (bill_1.eml, bill_2.eml).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_name: Name of the folder containing files (default: 'emails')
            - file_pattern: Base filename pattern (default: 'bill')
            - count: Expected number of files (default: 2)
            - dest: List of local destination filenames

    Returns:
        List of local filepaths where files were downloaded, or None for missing files
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    file_pattern = config.get('file_pattern', 'bill')
    count = config.get('count', 2)
    dest = config.get('dest', [f'pred_{i}.eml' for i in range(1, count + 1)])
    auth = GoogleAuth(settings_file=settings_file)
    drive = GoogleDrive(auth)
    folder_query = f"title = '{folder_name}' and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
    folder_list: GoogleDriveFileList = drive.ListFile({'q': folder_query}).GetList()
    if len(folder_list) == 0:
        logger.warning(f"Folder '{folder_name}' not found in Google Drive")
        return [None] * len(dest)
    folder_id = folder_list[0]['id']
    results = []
    for (i, dest_path) in enumerate(dest):
        file_num = i + 1
        filename = f'{file_pattern}_{file_num}.eml'
        file_query = f"title = '{filename}' and '{folder_id}' in parents"
        file_list: GoogleDriveFileList = drive.ListFile({'q': file_query}).GetList()
        if len(file_list) == 0:
            logger.warning(f"File '{filename}' not found in folder '{folder_name}'")
            results.append(None)
        else:
            file: GoogleDriveFile = file_list[0]
            try:
                file.GetContentFile(dest_path, mimetype=file['mimeType'])
                results.append(dest_path)
            except Exception as e:
                logger.error(f"Failed to download '{filename}': {e}")
                results.append(None)
    return results

def get_chrome_autofill_setting__cb0362f4(env, config):
    """
    Get the autofill enabled setting from Chrome preferences.

    Args:
        env: Environment object with controller and vm_platform
        config: Configuration dict (not used in this function)

    Returns:
        bool: True if autofill is enabled, False if disabled
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        autofill_enabled = data.get('autofill', {}).get('enabled', True)
        logger.info(f'Chrome autofill enabled setting: {autofill_enabled}')
        return autofill_enabled
    except Exception as e:
        logger.error(f'Error reading Chrome autofill setting: {e}')
        return True

def get_gdrive_file_check__cc11abb5(env, config: dict):
    """Check if paper_attachments folder exists on Google Drive and contains files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with settings_file for Google Drive

    Returns:
        dict: Dictionary with 'folder_exists' (bool), 'file_count' (int), and 'file_names' (list)
    """
    settings_file = config.get('settings_file', '')
    folder_check_cmd = f'''python3 -c "from desktop_env.evaluators.getters.chrome import get_googledrive_file; import json; folders = get_googledrive_file(None, {{'settings_file': '{settings_file}', 'query': \\"title='paper_attachments' and mimeType='application/vnd.google-apps.folder' and trashed=false\\"}}); print(json.dumps({{'exists': bool(folders), 'folder_id': folders[0]['id'] if folders else None}}))"'''
    folder_result = env.controller.run_bash_script(folder_check_cmd, timeout=30)
    folder_exists = False
    folder_id = None
    if folder_result and folder_result.get('returncode') == 0:
        try:
            import json
            output = folder_result.get('output', '').strip()
            folder_data = json.loads(output)
            folder_exists = folder_data.get('exists', False)
            folder_id = folder_data.get('folder_id')
        except (json.JSONDecodeError, ValueError, KeyError):
            folder_exists = False
            folder_id = None
    file_count = 0
    file_names = []
    if folder_exists and folder_id:
        files_check_cmd = f'''python3 -c "from desktop_env.evaluators.getters.chrome import get_googledrive_file; import json; files = get_googledrive_file(None, {{'settings_file': '{settings_file}', 'query': \\"'{folder_id}' in parents and trashed=false and mimeType!='application/vnd.google-apps.folder'\\"}}); print(json.dumps({{'count': len(files) if files else 0, 'names': [f['title'] for f in files] if files else []}}))"'''
        files_result = env.controller.run_bash_script(files_check_cmd, timeout=30)
        if files_result and files_result.get('returncode') == 0:
            try:
                import json
                output = files_result.get('output', '').strip()
                files_data = json.loads(output)
                file_count = files_data.get('count', 0)
                file_names = files_data.get('names', [])
            except (json.JSONDecodeError, ValueError, KeyError):
                file_count = 0
                file_names = []
    return {'folder_exists': folder_exists, 'file_count': file_count, 'file_names': file_names}

def get_extension_has_popup__36b763d0(env, config: Dict[str, Any]):
    """Check if extension has a popup action configured in manifest.json.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        bool: True if popup is configured, False otherwise
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    manifest_path = f'{extension_path}/manifest.json'
    try:
        content = env.controller.get_file(manifest_path)
        if content:
            manifest_data = json.loads(content)
            action = manifest_data.get('action', {})
            return 'default_popup' in action
        return False
    except Exception as e:
        logger.error(f'Error checking extension popup: {e}')
        return False

def get_gdrive_folder_contents__55cd0f01(env, config: dict):
    """Get contents of a Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with folder_name parameter

    Returns:
        Dict with folder info and file count
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        folder_name = config.get('folder_name', '')
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and 'root' in parents and trashed = false"
        folders = list_files(settings_file, folder_query)
        if not folders or not isinstance(folders, list) or len(folders) == 0:
            return {'folder_exists': False, 'file_count': 0}
        folder_id = folders[0].get('id', '')
        if not folder_id:
            return {'folder_exists': False, 'file_count': 0}
        files_query = f"'{folder_id}' in parents and mimeType = 'application/pdf' and trashed = false"
        files = list_files(settings_file, files_query)
        file_count = len(files) if isinstance(files, list) else 0
        return {'folder_exists': True, 'file_count': file_count, 'folder_id': folder_id}
    except Exception as e:
        print(f'Error getting Google Drive folder contents: {e}')
        return {'folder_exists': False, 'file_count': 0}

def get_chrome_standard_font_size__c9cdf92f31112c863794561a8e2fdc6b(env, config: Dict[str, str]):
    """
    Get the default font size setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Dictionary containing default_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        default_font_size = webprefs.get('default_font_size', 16)
        return {'default_font_size': default_font_size}
    except Exception as e:
        logger.error(f'Error getting default font size: {e}')
        return {'default_font_size': 16}

def get_chrome_default_font_size__8638fadde1ef14284e07488209d510be(env, config: Dict[str, str]):
    """
    Get the default font size setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Dictionary containing default_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        default_font_size = webprefs.get('default_font_size', 16)
        return {'default_font_size': default_font_size}
    except Exception as e:
        logger.error(f'Error getting default font size: {e}')
        return {'default_font_size': 16}

def get_googledrive_files_with_prefix__fb580d40(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in Google Drive folder matching a prefix"""
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        prefix = config.get('prefix', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            logger.warning(f'Folder not found with query: {folder_query}')
            return []
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        if prefix:
            matching_files = [f['title'] for f in file_list if f['title'].startswith(prefix)]
        else:
            matching_files = [f['title'] for f in file_list]
        logger.info(f"Found {len(matching_files)} files with prefix '{prefix}': {matching_files}")
        return sorted(matching_files)
    except Exception as e:
        logger.error(f'Error getting files from Google Drive: {e}')
        return []

def get_table_row_count_by_index__83874f16(env, config):
    """Get row count of specific table."""
    file_path = config.get('path')
    table_index = config.get('table_index', 0)
    if not file_path:
        return {'table_index': table_index, 'row_count': 0}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        return {'table_index': table_index, 'row_count': 0}
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        if table_index < len(doc.tables):
            row_count = len(doc.tables[table_index].rows)
        else:
            row_count = 0
        os.unlink(tmp_path)
        return {'table_index': table_index, 'row_count': row_count}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'table_index': table_index, 'row_count': 0}

def get_gdrive_file_metadata__6bd409d3(env, config: dict):
    """
    Get metadata about a file on Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'settings_file' and 'query'

    Returns:
        dict: File metadata including title, size, mimeType
    """
    from desktop_env.evaluators.getters.chrome import get_googledrive_file
    file_path = get_googledrive_file(env, config)
    if file_path is None:
        logger.warning('File not found on Google Drive')
        return None
    import os
    metadata = {'file_path': file_path, 'exists': os.path.exists(file_path) if file_path else False}
    if file_path and os.path.exists(file_path):
        metadata['size'] = os.path.getsize(file_path)
        logger.info(f"Retrieved file: {file_path}, size: {metadata['size']} bytes")
    return metadata

def get_chrome_show_home_button__aa9af3d4(env, config):
    """
    Get the show_home_button setting from Chrome browser preferences.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        bool: The value of browser.show_home_button setting (True if home button is shown, False if hidden)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        show_home_button = data.get('browser', {}).get('show_home_button', True)
        logger.info(f'Chrome show_home_button setting: {show_home_button}')
        return show_home_button
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return True

def get_chrome_font_size__82caef9e(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 3.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_chrome_font_size__fc8c3e93(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 4.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_gdrive_file_list__da0d777d0a1e3c8d735dbb81a2e45a5c(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a Google Drive folder by path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_path: List of folder names forming the path (e.g., ['TB_Export'])

    Returns:
        List of filenames found in the specified folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"'{parent_id}' in parents and title='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
            file_list = drive.ListFile({'q': query}).GetList()
            if not file_list:
                logger.warning(f"Folder '{folder_name}' not found in path {folder_path}")
                return []
            parent_id = file_list[0]['id']
        query = f"'{parent_id}' in parents and trashed=false"
        file_list = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in file_list]
        logger.info(f'Found {len(filenames)} files in Google Drive folder: {filenames}')
        return sorted(filenames)
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_gdrive_file_list__b0d15a73736689e26424d81f759f1528(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a Google Drive folder by path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_path: List of folder names forming the path (e.g., ['Bills_Backup'])

    Returns:
        List of filenames found in the specified folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"'{parent_id}' in parents and title='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
            file_list = drive.ListFile({'q': query}).GetList()
            if not file_list:
                logger.warning(f"Folder '{folder_name}' not found in path {folder_path}")
                return []
            parent_id = file_list[0]['id']
        query = f"'{parent_id}' in parents and trashed=false"
        file_list = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in file_list]
        logger.info(f'Found {len(filenames)} files in Google Drive folder: {filenames}')
        return sorted(filenames)
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_docx_table_info__818c616c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_default_web_browser__48359cc6c2268bcbb20cc6eebb7a0011(env, config: dict):
    """Gets the default web browser application on Linux.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: The most common default web browser application name
    """
    from desktop_env.evaluators.getters.general import get_vm_command_line
    os_type = env.vm_platform
    if os_type == 'Linux':
        mime_types = ['text/html', 'text/xml', 'application/xhtml+xml', 'x-scheme-handler/http', 'x-scheme-handler/https']
        apps = []
        for mime_type in mime_types:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', mime_type]})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_googledrive_files_list__d6fb1e53c50621e1a08efd7623119b0d(env, config):
    """Get list of files from Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Dict with 'settings_file' and 'folder_path' (list of folder names)

    Returns:
        List of filenames in the folder, or None if folder doesn't exist
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
    except ImportError:
        logger.error('pydrive not available')
        return None
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                logger.info(f"Folder '{folder_name}' not found")
                return []
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false and mimeType != 'application/vnd.google-apps.folder'"
        filelist = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in filelist]
        logger.info(f'Found files in folder: {filenames}')
        return filenames
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return None

def get_docx_table_content__044cff297ed58ab26ee69e56f58c3d19(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract table content from a DOCX file with section context.

    This getter reads a DOCX file and extracts all tables along with their
    section headings, enabling verification of table placement within document sections.

    Args:
        env: DesktopEnv instance with controller
        config: Configuration dict with 'path' key pointing to the DOCX file on VM

    Returns:
        Dictionary with tables data and section context, or None if error occurs
        Format: {
            'num_tables': int,
            'tables': [
                {
                    'num_rows': int,
                    'num_cols': int,
                    'data': [[cell_text, ...], ...],
                    'section': str  # The heading/section this table appears under
                },
                ...
            ]
        }
    """
    try:
        file_path = config.get('path', '')
        if not file_path:
            logger.error('No path specified in config')
            return None
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {file_path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            current_section = 'Document Start'
            tables_data = []
            table_index = 0
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    for para in doc.paragraphs:
                        if para._element == element:
                            if para.style.name.startswith('Heading') or para.style.name.startswith('Title'):
                                if para.text.strip():
                                    current_section = para.text.strip()
                                    logger.info(f'Found section heading: {current_section}')
                            break
                elif element.tag.endswith('tbl'):
                    if table_index < len(doc.tables):
                        table = doc.tables[table_index]
                        table_info = {'num_rows': len(table.rows), 'num_cols': len(table.columns) if table.rows else 0, 'data': [], 'section': current_section}
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            table_info['data'].append(row_data)
                        tables_data.append(table_info)
                        logger.info(f'Found table {table_index} in section: {current_section}')
                        table_index += 1
            result = {'num_tables': len(tables_data), 'tables': tables_data}
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting DOCX table content: {e}')
        return None

def get_recreation_cedarbreaks_html__ed948b8fac72e2c98dd5028aabd38bf3(env, config: Dict[str, Any]):
    """
    Get campsite data from recreation.gov page for Cedar Breaks search.
    Extracts campsite names, availability, and verifies Cedar Breaks location.
    This is a custom getter for the Cedar Breaks task variation.
    """
    logger.info(f'[RECREATION_CEDARBREAKS] Starting recreation.gov page processing for Cedar Breaks')
    logger.debug(f'[RECREATION_CEDARBREAKS] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 3
    timeout_ms = 60000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_CEDARBREAKS] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_CEDARBREAKS] Successfully connected to existing Chrome instance')
                except Exception as e:
                    logger.warning(f'[RECREATION_CEDARBREAKS] Failed to connect to existing Chrome instance: {e}')
                    logger.info(f'[RECREATION_CEDARBREAKS] Starting new Chrome instance...')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337', '--no-sandbox']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    logger.info(f"[RECREATION_CEDARBREAKS] Starting browser with command: {' '.join(command)}")
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup' + '/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_CEDARBREAKS] Successfully connected to new Chrome instance')
                if len(browser.contexts) == 0 or len(browser.contexts[0].pages) == 0:
                    logger.error(f'[RECREATION_CEDARBREAKS] No active pages found')
                    return None
                page = browser.contexts[0].pages[0]
                current_url = page.url
                logger.info(f'[RECREATION_CEDARBREAKS] Current URL: {current_url}')
                content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                result = {'cedar_breaks_found': False, 'has_campsite_results': False, 'campsites': [], 'closest_identified': False, 'url': current_url}
                page_text = soup.get_text().lower()
                if 'cedar breaks' in page_text or 'cedar-breaks' in page_text:
                    result['cedar_breaks_found'] = True
                    logger.info(f'[RECREATION_CEDARBREAKS] Cedar Breaks found in page content')
                campsite_cards = soup.find_all(['div', 'article'], class_=lambda x: x and ('campsite' in x.lower() or 'facility' in x.lower() or 'rec-card' in x.lower()))
                campsite_rows = soup.find_all('tr', class_=lambda x: x and ('campsite' in x.lower() or 'result' in x.lower()))
                campsite_items = soup.find_all(['li', 'div'], attrs={'data-facilityid': True})
                all_campsite_elements = campsite_cards + campsite_rows + campsite_items
                logger.info(f'[RECREATION_CEDARBREAKS] Found {len(all_campsite_elements)} potential campsite elements')
                if all_campsite_elements:
                    result['has_campsite_results'] = True
                    for (idx, element) in enumerate(all_campsite_elements[:10]):
                        campsite_info = {'position': idx, 'name': '', 'distance': '', 'available': False}
                        name_elem = element.find(['h3', 'h4', 'h5', 'a', 'span'], class_=lambda x: x and ('name' in x.lower() or 'title' in x.lower()))
                        if name_elem:
                            campsite_info['name'] = name_elem.get_text(strip=True)
                        elif element.get('data-facilityid'):
                            campsite_info['name'] = element.get_text(strip=True)[:100]
                        avail_text = element.get_text().lower()
                        if 'available' in avail_text or 'open' in avail_text:
                            campsite_info['available'] = True
                        elif 'sold out' in avail_text or 'unavailable' in avail_text or 'closed' in avail_text:
                            campsite_info['available'] = False
                        else:
                            campsite_info['available'] = True
                        distance_elem = element.find(text=lambda t: t and ('mile' in t.lower() or 'km' in t.lower() or 'mi' in t.lower()))
                        if distance_elem:
                            campsite_info['distance'] = distance_elem.strip()
                        result['campsites'].append(campsite_info)
                    if result['campsites']:
                        first_campsite = result['campsites'][0]
                        if first_campsite['available']:
                            result['closest_identified'] = True
                            logger.info(f"[RECREATION_CEDARBREAKS] Closest AVAILABLE campsite identified: {first_campsite['name']}")
                        else:
                            available_campsites = [c for c in result['campsites'] if c['available']]
                            if available_campsites:
                                result['closest_identified'] = True
                                logger.info(f"[RECREATION_CEDARBREAKS] Closest available campsite: {available_campsites[0]['name']} (position {available_campsites[0]['position']})")
                            else:
                                result['closest_identified'] = False
                                logger.warning(f'[RECREATION_CEDARBREAKS] No available campsites found')
                table_headers = soup.find_all(class_='camp-sortable-column-header')
                if len(table_headers) > 0:
                    result['has_campsite_results'] = True
                    logger.info(f'[RECREATION_CEDARBREAKS] Found {len(table_headers)} sortable table headers')
                logger.info(f"[RECREATION_CEDARBREAKS] Extraction complete. Cedar Breaks found: {result['cedar_breaks_found']}, Has results: {result['has_campsite_results']}, Campsites: {len(result['campsites'])}, Closest identified: {result['closest_identified']}")
                return result
        except Exception as e:
            logger.error(f'[RECREATION_CEDARBREAKS] Attempt {attempt + 1} failed: {e}')
            if attempt < max_retries - 1:
                logger.info(f'[RECREATION_CEDARBREAKS] Retrying in 2 seconds...')
                time.sleep(2)
            else:
                logger.error(f'[RECREATION_CEDARBREAKS] All retries exhausted')
                return None
    return None

def get_recreation_table_header__8356d858fc4a70c81a86864f34a14436(env, config: Dict[str, Any]):
    """
    Get table header from recreation.gov page to verify navigation.
    This is a simplified version that checks for specific table headers.
    """
    logger.info(f'[RECREATION_HEADER] Starting recreation.gov table header check')
    logger.debug(f'[RECREATION_HEADER] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 2
    timeout_ms = 30000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_HEADER] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_HEADER] Connected to Chrome')
                except Exception as e:
                    logger.warning(f'[RECREATION_HEADER] Failed to connect: {e}')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                page = browser.contexts[0].pages[0] if browser.contexts[0].pages else browser.contexts[0].new_page()
                current_url = page.url
                logger.info(f'[RECREATION_HEADER] Current URL: {current_url}')
                page.wait_for_load_state('networkidle', timeout=timeout_ms)
                html_content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                headers_elements = soup.find_all(class_='camp-sortable-column-header')
                result = {}
                for (idx, header) in enumerate(headers_elements):
                    result[str(idx)] = header.get_text(strip=True)
                logger.info(f'[RECREATION_HEADER] Found headers: {result}')
                return result
        except Exception as e:
            logger.error(f'[RECREATION_HEADER] Attempt {attempt + 1} failed: {e}')
            if attempt == max_retries - 1:
                logger.error(f'[RECREATION_HEADER] All attempts failed')
                return {}
            time.sleep(2)
    return {}

def get_chrome_ext_last_two__f117f7ab324ab80f0f8ad254a42cb210(env, config: Dict[str, str]) -> List[str]:
    """Get installed Chrome extension names.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of installed extension names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_name = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for id in all_extensions.keys():
            name = all_extensions[id]['manifest']['name']
            all_extensions_name.append(name)
        logger.info(f'Found installed extensions: {all_extensions_name}')
        return all_extensions_name
    except Exception as e:
        logger.error(f'Failed to get installed extensions: {e}')
        return []

def get_recreation_search_box__e0d56e1bf714f2c8f287d0502b986b63(env, config: Dict[str, Any]):
    """
    Check for the presence of a search box on recreation.gov page.
    """
    logger.info(f'[RECREATION_SEARCH] Starting recreation.gov search box check')
    logger.debug(f'[RECREATION_SEARCH] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 2
    timeout_ms = 30000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_SEARCH] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_SEARCH] Connected to Chrome')
                except Exception as e:
                    logger.warning(f'[RECREATION_SEARCH] Failed to connect: {e}')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                page = browser.contexts[0].pages[0] if browser.contexts[0].pages else browser.contexts[0].new_page()
                page.wait_for_load_state('networkidle', timeout=timeout_ms)
                html_content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html_content, 'html.parser')
                search_inputs = soup.find_all('input', attrs={'type': 'search'})
                search_inputs.extend(soup.find_all('input', attrs={'placeholder': lambda x: x and 'search' in x.lower()}))
                has_search_box = len(search_inputs) > 0
                logger.info(f'[RECREATION_SEARCH] Found {len(search_inputs)} search boxes')
                return {'has_search_box': has_search_box, 'count': len(search_inputs)}
        except Exception as e:
            logger.error(f'[RECREATION_SEARCH] Attempt {attempt + 1} failed: {e}')
            if attempt == max_retries - 1:
                logger.error(f'[RECREATION_SEARCH] All attempts failed')
                return {'has_search_box': False, 'count': 0}
            time.sleep(2)
    return {'has_search_box': False, 'count': 0}

def get_gdrive_eml_files_with_filenames__74b11cf6(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Query all .eml files in a Google Drive folder and return both content and filenames.

    This getter dynamically discovers all .eml files in the specified folder,
    downloads them, and returns both their content and actual filenames from Google Drive.
    This enables verification that files were named correctly based on their content.

    Args:
        env: Environment object
        config: Configuration dict with:
            - settings_file: path to Google Drive auth settings
            - folder_name: name of the folder to search (e.g., 'emails')
            - cache_dir_prefix: prefix for cache directory (default: 'eml_file')

    Returns:
        List of dicts, each containing:
            - 'filename': actual filename on Google Drive
            - 'content': bytes content of the file
            - 'local_path': local path where file was downloaded
        Returns None in list if file download fails.
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    cache_dir_prefix = config.get('cache_dir_prefix', 'eml_file')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false and 'root' in parents"
        folder_list: GoogleDriveFileList = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            logger.warning(f"Folder '{folder_name}' not found on Google Drive")
            return []
        folder: GoogleDriveFile = folder_list[0]
        folder_id = folder['id']
        eml_query = f"trashed = false and '{folder_id}' in parents"
        file_list: GoogleDriveFileList = drive.ListFile({'q': eml_query}).GetList()
        eml_files = [f for f in file_list if f['title'].endswith('.eml')]
        if len(eml_files) == 0:
            logger.warning(f"No .eml files found in folder '{folder_name}'")
            return []
        results = []
        for (idx, file) in enumerate(eml_files):
            try:
                filename = file['title']
                local_path = os.path.join(env.cache_dir, f'{cache_dir_prefix}_{idx}.eml')
                file.GetContentFile(local_path, mimetype=file['mimeType'])
                with open(local_path, 'rb') as f:
                    content = f.read()
                results.append({'filename': filename, 'content': content, 'local_path': local_path})
            except Exception as e:
                logger.error(f"Failed to download or read file '{file['title']}': {e}")
                results.append(None)
        return results
    except Exception as e:
        logger.error(f'Failed to query Google Drive: {e}')
        return []

def get_html_files_in_dir__98346f3b(env, config: dict):
    """Get list of HTML files in a directory with content snippets.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path to check

    Returns:
        Dict with:
            - count: Number of HTML files
            - files: List of dicts with filename and content info
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"""\npython3 -c "\nimport os\nimport glob\nimport json\nimport re\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path):\n    print(json.dumps({{'count': 0, 'files': []}}))\n    exit(0)\n\nhtml_files = glob.glob(os.path.join(dir_path, '*.html'))\nresult = {{'count': len(html_files), 'files': []}}\n\nfor file_path in html_files:\n    file_info = {{'filename': os.path.basename(file_path)}}\n    try:\n        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:\n            content = f.read(5000)  # Read first 5000 chars\n            # Extract title\n            title_match = re.search(r'<title[^>]*>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)\n            if title_match:\n                file_info['title'] = title_match.group(1).strip()\n            else:\n                file_info['title'] = ''\n            # Check for key markers\n            file_info['has_lilian_weng'] = 'Lilian Weng' in content or 'lilian weng' in content.lower()\n            file_info['has_agent_keyword'] = 'agent' in content.lower()\n            file_info['has_2023_06_23'] = '2023-06-23' in content\n            file_info['has_2024_02_05'] = '2024-02-05' in content\n            file_info['has_human_data_quality'] = 'human data quality' in content.lower()\n            file_info['content_snippet'] = content[:500]\n    except Exception as e:\n        file_info['error'] = str(e)\n\n    result['files'].append(file_info)\n\nprint(json.dumps(result))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to list HTML files: {result['error']}")
        return {'count': 0, 'files': []}
    import json
    try:
        data = json.loads(result['output'].strip())
        return data
    except:
        return {'count': 0, 'files': []}

def get_chrome_ext_subset__d7f8f8dc6935bc142e9b5dc629fdadfe(env, config: Dict[str, str]) -> List[str]:
    """Get installed Chrome extension names.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of installed extension names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_name = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for id in all_extensions.keys():
            name = all_extensions[id]['manifest']['name']
            all_extensions_name.append(name)
        logger.info(f'Found installed extensions: {all_extensions_name}')
        return all_extensions_name
    except Exception as e:
        logger.error(f'Failed to get installed extensions: {e}')
        return []

def get_gdrive_file_list__12a331f2(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_gdrive_file_list__839d38d5(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    from desktop_env.controllers.python import PythonController
    from typing import List
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_table_position__2d408821(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_docx_last_table_dims__3c609656(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_chrome_download_prompt__e5560b89(env, config: dict):
    """
    Check if Chrome's 'Ask where to save each file before downloading' setting is enabled.
    This checks the 'download.prompt_for_download' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if download prompt is enabled, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        prompt_for_download = data.get('download', {}).get('prompt_for_download', False)
        return 'true' if prompt_for_download else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_docx_table_count__8c6fecde(env, config: dict):
    """Count the number of tables in a docx file and extract document content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' parameters

    Returns:
        dict: Contains table count, table info, table content, and document text
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.error(f"Failed to get file from path: {config['path']}")
            return None
        import os
        cache_path = os.path.join(env.cache_dir, config['dest'])
        os.makedirs(os.path.dirname(cache_path), exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        doc = Document(cache_path)
        table_count = len(doc.tables)
        document_text = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                document_text.append(text)
        table_info = None
        table_content = None
        if table_count > 0:
            first_table = doc.tables[0]
            table_info = {'row_count': len(first_table.rows), 'col_count': len(first_table.columns)}
            table_cells = []
            for row in first_table.rows:
                row_cells = []
                for cell in row.cells:
                    row_cells.append(cell.text.strip())
                table_cells.append(row_cells)
            table_content = table_cells
        result = {'table_count': table_count, 'first_table_info': table_info, 'table_content': table_content, 'document_text': document_text}
        logger.info(f'Document data extracted: table_count={table_count}, has_text={len(document_text) > 0}')
        return result
    except Exception as e:
        logger.error(f'Error extracting document data: {e}')
        return None

def get_chrome_popups_blocked__6bbb0d83(env, config: dict):
    """
    Check if Chrome's pop-ups are blocked for all sites.
    This checks the 'profile.default_content_setting_values.popups' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if popups are blocked (value=2), "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        popups_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('popups', 0)
        return 'true' if popups_setting == 2 else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_table_position__dd5268aa(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_writer_table_first_n_rows__92370bc17bf2e2a5bd844c976eb32eb3(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract first N rows from Writer document table.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with table data
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            if len(doc.tables) == 0:
                return {'error': 'No tables found'}
            table = doc.tables[0]
            rows_data = []
            for row in table.rows:
                row_values = []
                for cell in row.cells:
                    text = cell.text.strip()
                    try:
                        if '.' in text:
                            val = float(text)
                        elif text.isdigit():
                            val = int(text)
                        else:
                            val = text
                        row_values.append(val)
                    except:
                        row_values.append(text)
                rows_data.append(row_values)
            return {'row_count': len(rows_data), 'rows': rows_data}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_gdrive_simplified_names__43be33c33f047d68d3b6cab5bd4d55ce(env, config: Dict[str, Any]) -> List[Optional[str]]:
    """Get email files from Google Drive with simplified filenames.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_name: Name of the folder containing files (default: 'emails')
            - filenames: List of expected simplified filenames
            - dest: List of local destination filenames

    Returns:
        List of local filepaths where files were downloaded, or None for missing files
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    filenames = config.get('filenames', [])
    dest = config.get('dest', [f'pred_{i}.eml' for i in range(len(filenames))])
    auth = GoogleAuth(settings_file=settings_file)
    drive = GoogleDrive(auth)
    folder_query = f"title = '{folder_name}' and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
    folder_list: GoogleDriveFileList = drive.ListFile({'q': folder_query}).GetList()
    if len(folder_list) == 0:
        logger.warning(f"Folder '{folder_name}' not found in Google Drive")
        return [None] * len(dest)
    folder_id = folder_list[0]['id']
    results = []
    for (filename, dest_path) in zip(filenames, dest):
        file_query = f"title = '{filename}' and '{folder_id}' in parents"
        file_list: GoogleDriveFileList = drive.ListFile({'q': file_query}).GetList()
        if len(file_list) == 0:
            logger.warning(f"File '{filename}' not found in folder '{folder_name}'")
            results.append(None)
        else:
            file: GoogleDriveFile = file_list[0]
            try:
                file.GetContentFile(dest_path, mimetype=file['mimeType'])
                results.append(dest_path)
            except Exception as e:
                logger.error(f"Failed to download '{filename}': {e}")
                results.append(None)
    return results

def get_gdrive_file_check__41049d6b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_gdrive_pdf_file__82f685a6(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_last_table_structure__e0b81f05(env, config):
    """Get last table structure with full content."""
    file_path = config.get('path')
    if not file_path:
        return {'row_count': 0, 'col_count': 0, 'headers': [], 'rows': []}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        return {'row_count': 0, 'col_count': 0, 'headers': [], 'rows': []}
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        if len(doc.tables) == 0:
            result = {'row_count': 0, 'col_count': 0, 'headers': [], 'rows': []}
        else:
            last_table = doc.tables[-1]
            row_count = len(last_table.rows)
            col_count = len(last_table.columns) if row_count > 0 else 0
            headers = [cell.text.strip() for cell in last_table.rows[0].cells] if row_count > 0 else []
            rows = []
            for i in range(1, row_count):
                row_data = [cell.text.strip() for cell in last_table.rows[i].cells]
                rows.append(row_data)
            result = {'row_count': row_count, 'col_count': col_count, 'headers': headers, 'rows': rows}
        os.unlink(tmp_path)
        return result
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'row_count': 0, 'col_count': 0, 'headers': [], 'rows': []}

def get_docx_last_table_dims__f8b20751(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_chrome_page_zoom__1c67d2f9(env, config: Dict[str, str]):
    """Get Chrome default page zoom level setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with default_zoom_level value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        zoom_level = data.get('partition', {}).get('default_zoom_level', {}).get('x', 0.0)
        return {'default_zoom_level': zoom_level}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'default_zoom_level': 0.0}

def get_bookmarks_in_folder__5641e4c0(env, config: dict):
    """Get list of bookmarks in a specific folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - folder_name: Name of bookmark folder to check

    Returns:
        List of dicts with 'title' and 'url' keys for each bookmark
    """
    folder_name = config.get('folder_name', 'Blog')
    command = f"""\npython3 -c "\nimport json\nimport os\n\n# Chrome bookmarks path\nbookmarks_path = os.path.expanduser('~/.config/google-chrome/Default/Bookmarks')\n\nif not os.path.exists(bookmarks_path):\n    print('[]')\n    exit(0)\n\nwith open(bookmarks_path, 'r') as f:\n    data = json.load(f)\n\ndef find_folder(node, target_name):\n    if node.get('type') == 'folder' and node.get('name') == target_name:\n        return node\n    for child in node.get('children', []):\n        result = find_folder(child, target_name)\n        if result:\n            return result\n    return None\n\n# Search in bookmark_bar and other\nfolder = None\nfor root_key in ['bookmark_bar', 'other']:\n    if root_key in data.get('roots', {{}}):\n        folder = find_folder(data['roots'][root_key], '{folder_name}')\n        if folder:\n            break\n\nif folder:\n    bookmarks = [{'title': item['name'], 'url': item['url']} for item in folder.get('children', []) if item.get('type') == 'url']\n    print(json.dumps(bookmarks))\nelse:\n    print('[]')\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to get bookmarks: {result['error']}")
        return []
    import json
    try:
        bookmarks = json.loads(result['output'].strip())
        return bookmarks
    except:
        return []

def get_csv_file_and_chrome_tab__29a47154(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if CSV file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict containing rules with csv_path

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    csv_path = config.get('rules', {}).get('csv_path')
    result = env.controller.get_file(csv_path)
    file_exists = result is not None and len(result) > 0
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    tabs_info = get_open_tabs_info(env, {})
    chrome_tabs = []
    if tabs_info:
        chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': csv_path}

def get_gdrive_text_file__953256df836603c8857d4495861e4b63(env, config: dict):
    """Get text file content from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'settings_file' and 'query_list'

    Returns:
        str: Content of the text file from Google Drive
    """
    from desktop_env.evaluators.getters.file import get_googledrive_file
    file_bytes = get_googledrive_file(env, config)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        print(f'Error decoding file: {e}')
        return ''

def get_chrome_font_size__88fb57a4(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 8.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_chrome_ext_single__abfa4b043befc7603aff32afef71bd1b(env, config: Dict[str, str]) -> List[str]:
    """Get installed Chrome extension names.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of installed extension names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_name = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for id in all_extensions.keys():
            name = all_extensions[id]['manifest']['name']
            all_extensions_name.append(name)
        logger.info(f'Found installed extensions: {all_extensions_name}')
        return all_extensions_name
    except Exception as e:
        logger.error(f'Failed to get installed extensions: {e}')
        return []

def get_chrome_third_party_cookies_blocked__15369290(env, config: dict):
    """
    Check if Chrome's third-party cookies are blocked.
    This checks the 'profile.block_third_party_cookies' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if third-party cookies are blocked, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        block_third_party = data.get('profile', {}).get('block_third_party_cookies', False)
        return 'true' if block_third_party else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_chrome_font_size__0ce3b14e(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 2.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_recreation_antelope_html__2aa2c046ed13ceae9537c9a8e566d558(env, config: Dict[str, Any]):
    """
    Get HTML content from recreation.gov page for Antelope Island search.
    Extracts location-specific information and availability data to verify the task was completed.
    This is a custom getter for the Antelope Island task variation.
    """
    logger.info(f'[RECREATION_ANTELOPE] Starting recreation.gov page processing for Antelope Island')
    logger.debug(f'[RECREATION_ANTELOPE] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 3
    timeout_ms = 60000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_ANTELOPE] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_ANTELOPE] Successfully connected to existing Chrome instance')
                except Exception as e:
                    logger.warning(f'[RECREATION_ANTELOPE] Failed to connect to existing Chrome instance: {e}')
                    logger.info(f'[RECREATION_ANTELOPE] Starting new Chrome instance...')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337', '--no-sandbox']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    logger.info(f"[RECREATION_ANTELOPE] Starting browser with command: {' '.join(command)}")
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup' + '/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_ANTELOPE] Successfully connected to new Chrome instance')
                if len(browser.contexts) == 0 or len(browser.contexts[0].pages) == 0:
                    logger.error(f'[RECREATION_ANTELOPE] No active pages found')
                    return None
                page = browser.contexts[0].pages[0]
                current_url = page.url
                logger.info(f'[RECREATION_ANTELOPE] Current URL: {current_url}')
                content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                result = {'location_verified': False, 'availability_found': False, 'url': current_url, 'search_term': None, 'first_available_date': None}
                url_lower = current_url.lower()
                if 'antelope' in url_lower or 'antelope-island' in url_lower:
                    result['location_verified'] = True
                    logger.info(f'[RECREATION_ANTELOPE] Location verified in URL: {current_url}')
                search_indicators = [soup.find('h1'), soup.find('title'), soup.find('input', {'type': 'search'}), soup.find(class_='rec-area-name'), soup.find(class_='facility-name')]
                for element in search_indicators:
                    if element:
                        text = element.get_text().strip() if hasattr(element, 'get_text') else str(element.get('value', ''))
                        if text and 'antelope' in text.lower():
                            result['search_term'] = text
                            result['location_verified'] = True
                            logger.info(f'[RECREATION_ANTELOPE] Found location reference in page: {text}')
                            break
                table_headers = soup.find_all(class_='camp-sortable-column-header')
                if len(table_headers) >= 2:
                    result['availability_found'] = True
                    logger.info(f'[RECREATION_ANTELOPE] Found availability table with {len(table_headers)} columns')
                date_selectors = [soup.find(class_='available'), soup.find(class_='rec-availability-date'), soup.find_all('td', class_='available')]
                for selector in date_selectors:
                    if selector:
                        if isinstance(selector, list) and len(selector) > 0:
                            date_text = selector[0].get_text().strip()
                        else:
                            date_text = selector.get_text().strip()
                        if date_text:
                            result['first_available_date'] = date_text
                            logger.info(f'[RECREATION_ANTELOPE] Found first available date: {date_text}')
                            break
                logger.info(f'[RECREATION_ANTELOPE] Extraction result: {result}')
                return result
        except Exception as e:
            logger.error(f'[RECREATION_ANTELOPE] Attempt {attempt + 1} failed: {e}')
            if attempt < max_retries - 1:
                logger.info(f'[RECREATION_ANTELOPE] Retrying in 2 seconds...')
                time.sleep(2)
            else:
                logger.error(f'[RECREATION_ANTELOPE] All retries exhausted')
                return None
    return None

def get_macys_url_parse__78b9ef3c3ae5dc84f1b8b31df8c5818e(env, config: Dict[str, str]):
    """
    Parse Macy's product URL for men's extra-large long-sleeve shirts with 30% discount.
    Variation 1: men's extra-large long-sleeve shirts with 30% discount.
    """
    result = {}
    active_tab_url = get_active_url_from_accessTree(env, config)
    if active_tab_url is None:
        return None
    parsed = urlparse(active_tab_url)
    path = unquote(parsed.path)
    result['mens_clothing'] = True if 'mens-clothing' in path else None
    path_parts = path.strip('/').split('/')
    key_value_json = {}
    shirts_flag = False
    short_sleeve_flag = False
    long_sleeve_flag = False
    if 'shirts' in path:
        shirts_flag = True
    if 'short-sleeve' in path:
        short_sleeve_flag = True
    if 'long-sleeve' in path:
        long_sleeve_flag = True
    for i in range(len(path_parts) - 1):
        if ',' in path_parts[i] and ',' in path_parts[i + 1]:
            keys = [k.strip() for k in path_parts[i].split(',')]
            values = [v.strip() for v in path_parts[i + 1].split(',')]
            for (k, v) in zip(keys, values):
                if k == 'Price_discount_range':
                    key_value_json[k] = [item.strip() for item in v.split('|')] if v else None
                else:
                    key_value_json[k] = v if v else None
                if k == 'Product_department' and (v == 'shirts' or v == 'Shirts' or v == 'Shirt'):
                    shirts_flag = True
                if k == 'Sleeve_length':
                    if v == 'short-sleeve' or v == 'Short Sleeve':
                        short_sleeve_flag = True
                    elif v == 'long-sleeve' or v == 'Long Sleeve':
                        long_sleeve_flag = True
            break
    for field in ['Men_regular_size_t', 'Price_discount_range', 'Sleeve_length']:
        if field not in key_value_json:
            key_value_json[field] = None
    result['shirts'] = shirts_flag if shirts_flag else None
    result['short_sleeve'] = short_sleeve_flag if short_sleeve_flag else None
    result['long_sleeve'] = long_sleeve_flag if long_sleeve_flag else None
    for key in config.get('parse_keys', []):
        if key in key_value_json:
            if key == 'Price_discount_range':
                if key_value_json[key] is not None:
                    if '50_PERCENT_ off & more' in key_value_json[key] and (not '30_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '50_PERCENT_ off & more'
                    elif '30_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '20_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '30_PERCENT_ off & more'
                    elif '20_PERCENT_ off & more' in key_value_json[key] and (not '50_PERCENT_ off & more' in key_value_json[key]) and (not '30_PERCENT_ off & more' in key_value_json[key]):
                        result[key] = '20_PERCENT_ off & more'
                    else:
                        result[key] = 'other_discount'
                else:
                    result[key] = 'no_discount'
            else:
                result[key] = key_value_json[key]
    return result

def get_docx_table_info__faf951c7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_chrome_setting_value__a785d845(env, config: Dict[str, str]):
    """
    Get the page zoom level setting from Chrome preferences.
    This getter extracts the default zoom level from Chrome's Preferences file.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        zoom_level = data.get('partition', {}).get('default_zoom_level', {}).get('x', 0.0)
        logger.info(f'[ZOOM_LEVEL] Retrieved zoom level: {zoom_level}')
        return {'zoom_level': zoom_level}
    except Exception as e:
        logger.error(f'Error getting zoom level: {e}')
        return {'zoom_level': 0.0}

def get_chrome_extension_manifest__6e0f538e8ca610f13ae0673161924889(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract Chrome extension manifest.json and check if required files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - manifest_path: Path to manifest.json on VM
            - required_files: List of file paths that should exist

    Returns:
        Dict with:
            - manifest: Parsed manifest.json content (dict)
            - files_exist: Dict mapping file paths to boolean (whether they exist)
            - all_files_exist: Boolean (True if all required files exist)
    """
    manifest_path = config.get('manifest_path', '')
    required_files = config.get('required_files', [])
    result = {'manifest': {}, 'files_exist': {}, 'all_files_exist': False}
    manifest_bytes = env.controller.get_file(manifest_path)
    if not manifest_bytes:
        return result
    try:
        manifest_content = manifest_bytes.decode('utf-8')
        result['manifest'] = json.loads(manifest_content)
    except (UnicodeDecodeError, json.JSONDecodeError):
        return result
    all_exist = True
    for file_path in required_files:
        file_bytes = env.controller.get_file(file_path)
        exists = file_bytes is not None and len(file_bytes) > 0
        result['files_exist'][file_path] = exists
        if not exists:
            all_exist = False
    result['all_files_exist'] = all_exist
    return result

def get_chrome_disable_fonts__df950818(env, config: Dict[str, str]):
    """Get Chrome disable custom fonts setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with disable_fonts_enabled value (boolean)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        disable_fonts = data.get('webkit', {}).get('webprefs', {}).get('fonts_disabled', False)
        return {'disable_fonts_enabled': disable_fonts}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'disable_fonts_enabled': False}

def get_gdrive_file_list__87199839(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_chrome_min_font_size__03b68be0b364e7ac1421db0dce73b670(env, config: Dict[str, str]):
    """
    Get the minimum font size setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Dictionary containing minimum_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        minimum_font_size = webprefs.get('minimum_font_size', 0)
        return {'minimum_font_size': minimum_font_size}
    except Exception as e:
        logger.error(f'Error getting minimum font size: {e}')
        return {'minimum_font_size': 0}

def get_gdrive_pdf_file__bf7218b9(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_chrome_ext_middle__484fc06534818b90d35ad638d0805c7c(env, config: Dict[str, str]) -> List[str]:
    """Get installed Chrome extension names.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of installed extension names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_name = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for id in all_extensions.keys():
            name = all_extensions[id]['manifest']['name']
            all_extensions_name.append(name)
        logger.info(f'Found installed extensions: {all_extensions_name}')
        return all_extensions_name
    except Exception as e:
        logger.error(f'Failed to get installed extensions: {e}')
        return []

def get_chrome_startup_url_removed__a85eebd24e563a97c17935bc46126aa1(env, config: Dict[str, str]):
    """
    Check if a specific URL has been removed from Chrome's startup URLs.
    Returns "true" if the URL is NOT in the startup list, "false" if it is.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        check_url = config.get('check_url', '')
        startup_urls = data.get('session', {}).get('startup_urls', [])
        logger.info(f'Current startup URLs: {startup_urls}')
        logger.info(f'Checking for URL: {check_url}')
        return 'false' if check_url in startup_urls else 'true'
    except Exception as e:
        logger.error(f'Error checking URL removal: {e}')
        return 'false'

def get_docx_table_structure__aa3c98fadda85fd2b11814a9c969fa8c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the structure of a table in a DOCX file, including dimensions and multiple cell values.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'table_index' keys

    Returns:
        Dict with table_exists, row_count, col_count, and cells data
    """
    try:
        vm_path = config['path']
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {vm_path}')
            return {'table_exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            table_index = config.get('table_index', 0)
            if table_index >= len(doc.tables):
                logger.error(f'Table index {table_index} out of range (total tables: {len(doc.tables)})')
                return {'table_exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}
            table = doc.tables[table_index]
            row_count = len(table.rows)
            col_count = len(table.columns)
            cells = {}
            for (row_idx, row) in enumerate(table.rows):
                for (col_idx, cell) in enumerate(row.cells):
                    cells[f'r{row_idx}c{col_idx}'] = cell.text.strip()
            return {'table_exists': True, 'row_count': row_count, 'col_count': col_count, 'cells': cells}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting table structure: {e}')
        return {'table_exists': False, 'row_count': 0, 'col_count': 0, 'cells': {}}

def get_chrome_font_size__25d0e193(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 5.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_vim_hlsearch_check__47d76eea6c988a4beb0cae6b4109cc24(env, config: dict):
    """
    Check if .vimrc file contains 'set hlsearch' configuration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        String output from the check script
    """
    script = '\nif [ -f ~/.vimrc ]; then\n    if grep -q "set hlsearch" ~/.vimrc; then\n        echo "The File Has Set Hlsearch!"\n    else\n        echo "The File Does Not Have Set Hlsearch!"\n    fi\nelse\n    echo "The .vimrc File Does Not Exist!"\nfi\n'
    result = env.controller.run_bash_script(script, timeout=10)
    output = result.get('output', '').strip()
    logger.info(f'Vim hlsearch check result: {output}')
    return output

def get_gdrive_eml_files__26d2516b888edf7c1b328cca7acaf9b7(env, config: Dict[str, Any]) -> List[str]:
    """Get list of .eml files in Google Drive 'emails' folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_name: Folder name to check (default: 'emails')

    Returns:
        List of .eml filenames found in the folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'emails')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and 'root' in parents and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            logger.info(f"Folder '{folder_name}' not found")
            return []
        folder_id = folder_list[0]['id']
        file_query = f"'{folder_id}' in parents and trashed = false and title contains '.eml'"
        file_list = drive.ListFile({'q': file_query}).GetList()
        eml_files = [f['title'] for f in file_list if f['title'].endswith('.eml')]
        logger.info(f'Found {len(eml_files)} .eml files: {eml_files}')
        return sorted(eml_files)
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_chrome_third_party_cookies__1394774d(env, config: Dict[str, str]):
    """Get Chrome third-party cookies blocking setting from preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Third-party cookies settings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        block_third_party = data.get('profile', {}).get('block_third_party_cookies', False)
        return {'block_third_party': block_third_party}
    except Exception as e:
        logger.error(f'Error getting third-party cookies setting: {e}')
        return {'block_third_party': False}

def get_gdrive_filenames__71c23132811d122bc61dca33636b8f81(env, config: Dict[str, Any]) -> List[str]:
    """
    Get the list of file names in a specific Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to check (default: 'figures')

    Returns:
        List[str]: Sorted list of file names in the folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'figures')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and trashed = false and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if not folder_list:
            return []
        folder_id = folder_list[0]['id']
        file_query = f"trashed = false and '{folder_id}' in parents"
        file_list = drive.ListFile({'q': file_query}).GetList()
        filenames = sorted([f['title'] for f in file_list])
        return filenames
    except Exception as e:
        import logging
        logger = logging.getLogger('desktopenv.getter.googledrive')
        logger.error(f'Error getting Google Drive file names: {e}')
        return []

def get_docx_table_info__f2fed625(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_gdrive_file_list__45753efe(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_gdrive_file_check__b2131ee6(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_writer_table_rows__b3ed6b27f25c4e432395f7240c942359(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract table rows from a LibreOffice Writer document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the docx file

    Returns:
        Dict with table data including row count and cell values
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'rows': 0, 'error': 'File not found'}
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            if len(doc.tables) == 0:
                return {'rows': 0, 'error': 'No tables found'}
            table = doc.tables[0]
            rows_data = []
            for row in table.rows:
                row_values = []
                for cell in row.cells:
                    text = cell.text.strip()
                    try:
                        if '.' in text:
                            val = float(text)
                        else:
                            val = int(text)
                        row_values.append(val)
                    except:
                        row_values.append(text)
                rows_data.append(row_values)
            return {'rows': len(rows_data), 'cols': len(rows_data[0]) if rows_data else 0, 'data': rows_data}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'rows': 0, 'error': str(e)}

def get_gdrive_pdf_file__c340b25e(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_table_count_and_first_table__638d8a63(env, config: Dict[str, Any]):
    """Get the count of tables and content of first table.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path

    Returns:
        dict: {'table_count': int, 'first_table_data': list of lists}
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return {'table_count': 0, 'first_table_data': []}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        logger.error(f'Could not read file: {file_path}')
        return {'table_count': 0, 'first_table_data': []}
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
            tmp_file.write(file_data)
            tmp_path = tmp_file.name
        doc = Document(tmp_path)
        table_count = len(doc.tables)
        first_table_data = []
        if table_count > 0:
            first_table = doc.tables[0]
            for row in first_table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                first_table_data.append(row_data)
        os.unlink(tmp_path)
        return {'table_count': table_count, 'first_table_data': first_table_data}
    except Exception as e:
        logger.error(f'Error reading document: {e}')
        return {'table_count': 0, 'first_table_data': []}

def get_docx_table_info__bf2c005c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_docx_table_structure__deda270825a0b396cee34e9436d907d9(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the structure (rows, columns) of a specific table in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'table_index' keys

    Returns:
        Dict with 'rows' and 'columns' keys, or empty dict on error
    """
    try:
        vm_path = config['path']
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {vm_path}')
            return {}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            table_index = config.get('table_index', 0)
            if table_index >= len(doc.tables):
                logger.error(f'Table index {table_index} out of range (total tables: {len(doc.tables)})')
                return {}
            table = doc.tables[table_index]
            return {'rows': len(table.rows), 'columns': len(table.columns)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting table structure: {e}')
        return {}

def get_chrome_font_size__baaf192c(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 1.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_gdrive_file_exists__3ce93f6ce10147dcc2489b34874a9f3b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on Google Drive and get its metadata.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - settings_file: Path to Google Drive settings file
            - query: Query string to find the file

    Returns:
        Dict with keys:
            - exists: bool - whether the file exists
            - file_count: int - number of matching files found
            - file_name: str or None - name of the first matching file
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
    except ImportError:
        return {'exists': False, 'file_count': 0, 'file_name': None}
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    query = config.get('query', '')
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        filelist = drive.ListFile({'q': query}).GetList()
        if len(filelist) > 0:
            return {'exists': True, 'file_count': len(filelist), 'file_name': filelist[0]['title']}
        else:
            return {'exists': False, 'file_count': 0, 'file_name': None}
    except Exception as e:
        return {'exists': False, 'file_count': 0, 'file_name': None, 'error': str(e)}

def get_webext_dir__951cee96(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_chrome_experiments_partial__70b32da51f968d0aa45e836eecc52bdb(env, config: Dict[str, str]):
    """
    Get enabled Chrome experiments and return them as a list.
    This getter is used to check partial matches of enabled experiments.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, but required by framework)

    Returns:
        List of enabled experiment names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Local State'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enabled_labs_experiments = data.get('browser', {}).get('enabled_labs_experiments', [])
        experiment_names = [exp.split('@')[0] for exp in enabled_labs_experiments]
        return experiment_names
    except Exception as e:
        logger.error(f'Error getting enabled experiments: {e}')
        return []

def get_gdrive_file_check__94eb0e23(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_chrome_min_font_size__a5307030(env, config: Dict[str, str]):
    """Get Chrome minimum font size setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with minimum_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {'minimum_font_size': 0})
        return {'minimum_font_size': webprefs.get('minimum_font_size', 0)}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'minimum_font_size': 0}

def get_gdrive_file_list__f478ac41(env, config: dict):
    """Get list of files in Google Drive root folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with query parameter

    Returns:
        List of file names found in Google Drive
    """
    import sys
    import os
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../../../../src/envs/osworld_env/desktop_env/evaluators/getters'))
    try:
        from googledrive import list_files
        settings_file = config.get('settings_file', '')
        query = config.get('query', '')
        files = list_files(settings_file, query)
        if isinstance(files, list):
            file_names = [f.get('name', '') for f in files if isinstance(f, dict)]
            return sorted(file_names)
        return []
    except Exception as e:
        print(f'Error getting Google Drive file list: {e}')
        return []

def get_git_remote_url__f8582c17(env, config: dict):
    """
    Get the remote URL of a git repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        str: Remote URL or empty string
    """
    repo_path = config.get('repo_path', '')
    command = f'cd "{repo_path}" && git remote get-url origin 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        return result['output'].strip()
    return ''

def get_docx_last_table_dims__86f8b2c6(env, config: Dict[str, Any]):
    """Get table count and dimensions of the last table in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'table_count', 'last_table_dims' (rows/cols), and 'is_empty' keys
              Returns None if file cannot be read
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        table_count = len(doc.tables)
        if table_count == 0:
            return {'table_count': 0, 'last_table_dims': None, 'is_empty': None}
        last_table = doc.tables[-1]
        rows = len(last_table.rows)
        cols = len(last_table.columns)
        is_empty = True
        for row in last_table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    is_empty = False
                    break
            if not is_empty:
                break
        return {'table_count': table_count, 'last_table_dims': {'rows': rows, 'cols': cols}, 'is_empty': is_empty}
    except Exception as e:
        return None

def get_gdrive_file_check__9b0aebe9(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_extension_manifest_version__fb0f6c5b(env, config):
    """
    Get the manifest_version field from the 'Hello Extensions' extension.

    Args:
        env: Environment object
        config: Configuration dict with 'extension_name' key

    Returns:
        int: The manifest_version value from the extension's manifest, or None if not found
    """
    extension_name = config.get('extension_name', 'Hello Extensions')
    os_type = env.vm_platform
    logger.info(f'[EXTENSION_MANIFEST_VERSION] Looking for extension: {extension_name}')
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        logger.error(f'[EXTENSION_MANIFEST_VERSION] Unsupported operating system: {os_type}')
        return None
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        logger.info(f'[EXTENSION_MANIFEST_VERSION] Successfully read Preferences file')
        all_extensions = data.get('extensions', {}).get('settings', {})
        logger.info(f'[EXTENSION_MANIFEST_VERSION] Found {len(all_extensions)} installed extensions')
        for (extension_id, extension_data) in all_extensions.items():
            manifest = extension_data.get('manifest', {})
            name = manifest.get('name', '')
            logger.debug(f'[EXTENSION_MANIFEST_VERSION] Checking extension: {name} (ID: {extension_id})')
            if name == extension_name:
                manifest_version = manifest.get('manifest_version')
                logger.info(f"[EXTENSION_MANIFEST_VERSION] Found extension '{extension_name}' with manifest_version: {manifest_version}")
                return manifest_version
        logger.warning(f"[EXTENSION_MANIFEST_VERSION] Extension '{extension_name}' not found in installed extensions")
        logger.info(f"[EXTENSION_MANIFEST_VERSION] Available extensions: {[data.get('manifest', {}).get('name', '') for data in all_extensions.values()]}")
        return None
    except Exception as e:
        logger.error(f'[EXTENSION_MANIFEST_VERSION] Error reading Chrome Preferences file: {e}')
        return None

def get_chrome_min_font_size__986d3358722e45a6a9be7a39f21f689e(env, config: Dict[str, str]):
    """Get Chrome's minimum font size setting from Preferences file."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {'minimum_font_size': 0})
        return webprefs
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'minimum_font_size': 0}

def get_googledrive_files_list__ba61b137046435f47239f8911466a875(env, config):
    """Get list of files from Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Dict with 'settings_file' and 'folder_path' (list of folder names)

    Returns:
        List of filenames in the folder, or None if folder doesn't exist
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
    except ImportError:
        logger.error('pydrive not available')
        return None
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                logger.info(f"Folder '{folder_name}' not found")
                return []
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false and mimeType != 'application/vnd.google-apps.folder'"
        filelist = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in filelist]
        logger.info(f'Found files in folder: {filenames}')
        return filenames
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return None

def get_extension_manifest__b5fa3477caae1bab3fd0d8b19ef4dfc7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract manifest.json content from browser extension directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to manifest.json

    Returns:
        Dictionary containing manifest.json content, or empty dict if file not found
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest.json from VM: {path}')
            return {}
        manifest_data = json.loads(file_bytes.decode('utf-8'))
        return manifest_data
    except json.JSONDecodeError as e:
        logger.error(f'JSON decode error in manifest.json: {e}')
        return {}
    except Exception as e:
        logger.error(f'Error reading manifest.json from {path}: {e}')
        return {}

def get_gdrive_nested_files__688ade84b6705b21f2a27dc4863ba216(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in nested Google Drive folder structure.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - file_paths: List of file paths, each path is a list of folder/file names

    Returns:
        List of file paths that exist (formatted as "folder/subfolder/file.ext")
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    file_paths = config.get('file_paths', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        found_files = []
        for path_components in file_paths:
            parent_id = 'root'
            path_exists = True
            for (i, component) in enumerate(path_components):
                is_last = i == len(path_components) - 1
                if is_last:
                    query = f"'{parent_id}' in parents and title='{component}' and trashed=false"
                else:
                    query = f"'{parent_id}' in parents and title='{component}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
                file_list = drive.ListFile({'q': query}).GetList()
                if not file_list:
                    logger.warning(f"Component '{component}' not found in path {path_components}")
                    path_exists = False
                    break
                if not is_last:
                    parent_id = file_list[0]['id']
            if path_exists:
                path_str = '/'.join(path_components)
                found_files.append(path_str)
                logger.info(f'Found file at path: {path_str}')
        return found_files
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_chrome_fixed_font_family__a0c1ef41(env, config: Dict[str, str]):
    """Get Chrome fixed (monospace) font family setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with fixed_font_family value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        fonts = data.get('webkit', {}).get('webprefs', {}).get('fonts', {}).get('fixed', {})
        fixed_font = fonts.get('Zyyy', 'Courier New')
        return {'fixed_font_family': fixed_font}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'fixed_font_family': 'Courier New'}

def get_chrome_password_breach_alerts__2107d226(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to alert about password breaches.
    Returns 'true' if password breach detection is enabled, 'false' otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        profile = data.get('profile', {})
        leak_detection_enabled = profile.get('password_manager_leak_detection', False)
        return 'true' if leak_detection_enabled else 'false'
    except Exception as e:
        logger.error(f'Error checking Chrome password breach alerts: {e}')
        return 'false'

def get_clear_cookies_on_exit__936eec5ac877849802aeaa5f0c43e44f(env, config: Dict[str, str]):
    """
    Check if Chrome is configured to clear cookies and site data on exit.

    This checks the actual cookie clearing preferences in Chrome's Preferences file.
    The setting is typically found in:
    - browser.clear_data.cookies_basic (for newer Chrome versions)
    - profile.exit_type with value "Normal"
    - browser.clear_data settings

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        str: "true" if cookies are set to clear on exit, "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        return 'false'
    try:
        content = env.controller.get_file(preference_file_path)
        if not content:
            return 'false'
        data = json.loads(content)
        if 'browser' in data and 'clear_data' in data['browser']:
            clear_data = data['browser']['clear_data']
            if clear_data.get('cookies_basic', False):
                return 'true'
        if 'profile' in data:
            profile = data['profile']
            if 'content_settings' in profile and 'exceptions' in profile['content_settings']:
                exceptions = profile['content_settings']['exceptions']
                if 'cookies' in exceptions:
                    cookies = exceptions['cookies']
                    for (pattern, settings) in cookies.items():
                        if isinstance(settings, dict) and settings.get('setting') == 4:
                            return 'true'
            if 'default_content_setting_values' in profile:
                dcv = profile['default_content_setting_values']
                if dcv.get('cookies') == 4:
                    return 'true'
        if 'local_state' in data:
            local_state = data['local_state']
            if 'clear_browsing_data' in local_state:
                cbd = local_state['clear_browsing_data']
                if cbd.get('cookies', False):
                    return 'true'
        return 'false'
    except Exception as e:
        return 'false'

def get_webext_dir__81e9cd46(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_docx_table_content__0da226e8b4b1ba563f27816711980e35(env, config: Dict[str, Any]) -> List[List[str]]:
    """Extract table content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        List of table rows, where each row is a list of cell values
    """
    file_path = config.get('path')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_tables.append(table_data)
        return all_tables
    finally:
        os.unlink(tmp_path)

def get_chrome_zoom_level__d32e355a(env, config):
    """Get Chrome's default zoom level setting from Preferences file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Zoom level settings including 'default_zoom_level'
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        partition = data.get('partition', {})
        default_zoom_level = partition.get('default_zoom_level', {})
        zoom_level = default_zoom_level.get('x', 0.0)
        return {'default_zoom_level': zoom_level}
    except Exception as e:
        logger.error(f'Error getting zoom level: {e}')
        return {'default_zoom_level': 0.0}

def get_webext_manifest__4b1e45b68c5d85573f69177601102dcb(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract manifest.json from a web extension project directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to manifest.json on VM)

    Returns:
        Dict containing the manifest JSON, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest file from VM: {path}')
            return None
        manifest = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded manifest from {path}')
        return manifest
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading manifest from {path}: {e}')
        return None

def get_chrome_ext_productivity__3ce2c199effd9eefc89344b99f4cd5a7(env, config: Dict[str, str]) -> List[str]:
    """Get installed Chrome extension names.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        List of installed extension names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_name = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for id in all_extensions.keys():
            name = all_extensions[id]['manifest']['name']
            all_extensions_name.append(name)
        logger.info(f'Found installed extensions: {all_extensions_name}')
        return all_extensions_name
    except Exception as e:
        logger.error(f'Failed to get installed extensions: {e}')
        return []

def get_docx_table_data__54b06cb187414c0c4afc1326d01127a0(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract table data and document structure from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the VM file path

    Returns:
        Dictionary containing:
        - 'tables': List of tables, where each table is a list of rows, and each row is a list of cell values
        - 'structure': List of document elements in order, with type ('paragraph' or 'table') and content/index
        Returns None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        structure = []
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cell_values = [cell.text.strip() for cell in row.cells]
                table_rows.append(cell_values)
            tables_data.append(table_rows)
        table_index = 0
        for element in doc.element.body:
            if element.tag.endswith('p'):
                for para in doc.paragraphs:
                    if para._element == element:
                        structure.append({'type': 'paragraph', 'text': para.text.strip(), 'style': para.style.name if para.style else None})
                        break
            elif element.tag.endswith('tbl'):
                structure.append({'type': 'table', 'index': table_index})
                table_index += 1
        return {'tables': tables_data, 'structure': structure}
    except Exception as e:
        print(f'Error reading docx: {e}')
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_docx_table_content__1c032efed907d41d75501f3895adca32(env, config: Dict[str, Any]) -> List[List[str]]:
    """Extract table content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        List of table rows, where each row is a list of cell values
    """
    file_path = config.get('path')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_tables.append(table_data)
        return all_tables
    finally:
        os.unlink(tmp_path)

def get_chrome_password_manager__09e00537(env, config):
    """
    Get the password manager enabled state from Chrome Preferences.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        bool: True if password_manager_enabled is true, False otherwise
    """
    os_type = env.vm_platform
    logger.info('[PASSWORD_MANAGER_GETTER] Getting Chrome password manager state')
    logger.info(f'[PASSWORD_MANAGER_GETTER] OS type: {os_type}')
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        logger.error(f'[PASSWORD_MANAGER_GETTER] Unsupported operating system: {os_type}')
        raise Exception('Unsupported operating system')
    logger.info(f'[PASSWORD_MANAGER_GETTER] Preference file path: {preference_file_path}')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        password_manager_enabled = data.get('profile', {}).get('password_manager_enabled', True)
        logger.info(f'[PASSWORD_MANAGER_GETTER] Password manager enabled: {password_manager_enabled}')
        return password_manager_enabled
    except Exception as e:
        logger.error(f'[PASSWORD_MANAGER_GETTER] Error reading Chrome Preferences: {e}')
        return True

def get_gdrive_file_check__ada1e84e(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_gdrive_pdf_file__7e0bdf48(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_extensions_page_content__ae6416e4(env, config: Dict[str, Any]) -> str:
    """Navigate to chrome://extensions and get page content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        HTML content from chrome://extensions page as string
    """
    try:
        env.controller.get_activate_window_command('Google Chrome')
        env.controller.execute_pyautogui_command("pyautogui.hotkey('ctrl', 'l'); import time; time.sleep(0.3)")
        env.controller.execute_pyautogui_command("pyautogui.write('chrome://extensions', interval=0.05)")
        env.controller.execute_pyautogui_command("pyautogui.press('enter'); import time; time.sleep(2.0)")
        max_retries = 3
        for attempt in range(max_retries):
            try:
                result = env.controller.execute_python_command('\nimport subprocess\nimport json\nimport time\n\n# Get all Chrome tabs\nresult = subprocess.run([\'curl\', \'-s\', \'http://localhost:9222/json\'], capture_output=True, text=True)\ntabs = json.loads(result.stdout)\n\nfor tab in tabs:\n    if \'chrome://extensions\' in tab.get(\'url\', \'\'):\n        ws_url = tab[\'webSocketDebuggerUrl\']\n        # Get page content via CDP\n        try:\n            import websocket\n        except ImportError:\n            print("ERROR: websocket library not available")\n            break\n\n        try:\n            ws = websocket.create_connection(ws_url, timeout=5)\n            ws.send(json.dumps({"id": 1, "method": "Runtime.evaluate", "params": {"expression": "document.body.innerText"}}))\n            response = json.loads(ws.recv())\n            ws.close()\n            if \'result\' in response and \'result\' in response[\'result\']:\n                print(response[\'result\'][\'result\'].get(\'value\', \'\'))\n                break\n        except Exception as e:\n            print(f"ERROR connecting to websocket: {e}")\n            break\n')
                content = result.get('output', '')
                if content and (not content.startswith('ERROR')):
                    return content
                if attempt < max_retries - 1:
                    import time
                    time.sleep(1.0)
            except Exception as inner_e:
                logger.warning(f'Retry {attempt + 1}/{max_retries} failed: {inner_e}')
                if attempt < max_retries - 1:
                    import time
                    time.sleep(1.0)
        return result.get('output', '') if 'result' in locals() else ''
    except Exception as e:
        logger.error(f'Error getting chrome://extensions content: {e}')
        return ''

def get_extension_manifest_version__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """Get manifest version and unpacked status for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict mapping extension names to dicts with 'manifest_version' and 'is_unpacked' keys
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    max_retries = 3
    retry_delay = 2
    for attempt in range(max_retries):
        try:
            file_exists_result = env.controller.execute_python_command(f"import os; print(os.path.exists('{preference_file_path}'))")
            file_exists = file_exists_result['output'].strip().lower() == 'true'
            if not file_exists:
                if attempt < max_retries - 1:
                    logger.warning(f'Preferences file not found, attempt {attempt + 1}/{max_retries}. Retrying...')
                    time.sleep(retry_delay)
                    continue
                else:
                    logger.error(f'Preferences file not found after {max_retries} attempts')
                    return {}
            content = env.controller.get_file(preference_file_path)
            data = json.loads(content)
            extensions_info = {}
            all_extensions = data.get('extensions', {}).get('settings', {})
            if not all_extensions and attempt < max_retries - 1:
                logger.warning(f'Extensions settings empty, attempt {attempt + 1}/{max_retries}. Retrying...')
                time.sleep(retry_delay)
                continue
            for ext_id in all_extensions.keys():
                ext_data = all_extensions[ext_id]
                name = ext_data.get('manifest', {}).get('name', '')
                manifest_version = ext_data.get('manifest', {}).get('manifest_version', 0)
                path = ext_data.get('path', '')
                is_unpacked = bool(path and (not path.startswith('chrome://')))
                if name:
                    extensions_info[name] = {'manifest_version': manifest_version, 'is_unpacked': is_unpacked}
            return extensions_info
        except json.JSONDecodeError as e:
            if attempt < max_retries - 1:
                logger.warning(f'Error parsing Preferences JSON, attempt {attempt + 1}/{max_retries}. Retrying... Error: {e}')
                time.sleep(retry_delay)
                continue
            else:
                logger.error(f'Error parsing Preferences JSON after {max_retries} attempts: {e}')
                return {}
        except Exception as e:
            if attempt < max_retries - 1:
                logger.warning(f'Error reading Chrome Preferences file, attempt {attempt + 1}/{max_retries}. Retrying... Error: {e}')
                time.sleep(retry_delay)
                continue
            else:
                logger.error(f'Error reading Chrome Preferences file for extension manifest versions after {max_retries} attempts: {e}')
                return {}
    return {}

def get_docx_table_bold_status__67405416(env, config):
    """Check if all table cells are bold."""
    file_path = config.get('path', '/home/user/Desktop/presentation_instruction_2023_Feb.docx')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'all_bold': False}
    doc = Document(io.BytesIO(file_bytes))
    if len(doc.tables) == 0:
        return {'all_bold': False}
    table = doc.tables[0]
    all_bold = True
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                for run in para.runs:
                    if run.text.strip() and (not run.bold):
                        all_bold = False
                        break
                if not all_bold:
                    break
            if not all_bold:
                break
        if not all_bold:
            break
    return {'all_bold': all_bold}

def get_vm_file(env, config: Dict[str, Any]):
    """Get file from VM - simplified version for this getter."""
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        return cache_path
    except Exception as e:
        logger.error(f'Error getting file {path}: {e}')
        return None

def get_gdrive_pdf_file__a70aba9d(env, config: Dict[str, Any]) -> Any:
    """Get PDF file from Google Drive folder.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to search in
            - file_patterns: List of possible PDF filenames
            
    Returns:
        str: Path to downloaded PDF file, or None if not found
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'forms')
    file_patterns = config.get('file_patterns', ['form.pdf', 'form.docx.pdf'])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return None
        folder_id = folder_list[0]['id']
        for pattern in file_patterns:
            file_query = f"(title = '{pattern}') and '{folder_id}' in parents and trashed = false"
            file_list = drive.ListFile({'q': file_query}).GetList()
            if len(file_list) > 0:
                pdf_file = file_list[0]
                dest_path = os.path.join(env.cache_dir, 'form.pdf')
                pdf_file.GetContentFile(dest_path, mimetype=pdf_file['mimeType'])
                return dest_path
        return None
    except Exception as e:
        print(f'Error retrieving PDF from Google Drive: {e}')
        return None

def get_docx_table_info__2aa481b8(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_chrome_extension_manifest__70453fbf044a41274b6b7dc7e909fb11(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract Chrome extension manifest.json and check if required files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - manifest_path: Path to manifest.json on VM
            - required_files: List of file paths that should exist

    Returns:
        Dict with:
            - manifest: Parsed manifest.json content (dict)
            - files_exist: Dict mapping file paths to boolean (whether they exist)
            - all_files_exist: Boolean (True if all required files exist)
    """
    manifest_path = config.get('manifest_path', '')
    required_files = config.get('required_files', [])
    result = {'manifest': {}, 'files_exist': {}, 'all_files_exist': False}
    manifest_bytes = env.controller.get_file(manifest_path)
    if not manifest_bytes:
        return result
    try:
        manifest_content = manifest_bytes.decode('utf-8')
        result['manifest'] = json.loads(manifest_content)
    except (UnicodeDecodeError, json.JSONDecodeError):
        return result
    all_exist = True
    for file_path in required_files:
        file_bytes = env.controller.get_file(file_path)
        exists = file_bytes is not None and len(file_bytes) > 0
        result['files_exist'][file_path] = exists
        if not exists:
            all_exist = False
    result['all_files_exist'] = all_exist
    return result

def get_chrome_extension_manifest__61167ddb747c3a88a31446374f5a958f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract Chrome extension manifest.json and check if required files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters:
            - manifest_path: Path to manifest.json on VM
            - required_files: List of file paths that should exist

    Returns:
        Dict with:
            - manifest: Parsed manifest.json content (dict)
            - files_exist: Dict mapping file paths to boolean (whether they exist)
            - all_files_exist: Boolean (True if all required files exist)
    """
    manifest_path = config.get('manifest_path', '')
    required_files = config.get('required_files', [])
    result = {'manifest': {}, 'files_exist': {}, 'all_files_exist': False}
    manifest_bytes = env.controller.get_file(manifest_path)
    if not manifest_bytes:
        return result
    try:
        manifest_content = manifest_bytes.decode('utf-8')
        result['manifest'] = json.loads(manifest_content)
    except (UnicodeDecodeError, json.JSONDecodeError):
        return result
    all_exist = True
    for file_path in required_files:
        file_bytes = env.controller.get_file(file_path)
        exists = file_bytes is not None and len(file_bytes) > 0
        result['files_exist'][file_path] = exists
        if not exists:
            all_exist = False
    result['all_files_exist'] = all_exist
    return result

def get_block_third_party_cookies__d0393d13df595b6db99860dc4f30ea7b(env, config: Dict[str, str]):
    """Get Chrome's third-party cookie blocking setting from Preferences file."""
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        block_third_party = data.get('profile', {}).get('block_third_party_cookies', False)
        return {'block_third_party_cookies': block_third_party}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'block_third_party_cookies': False}

def get_googledrive_folder_details__68f5fe6b(env, config: Dict[str, Any]):
    """Get detailed info about files in a Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with folder_path

    Returns:
        dict: {"file_count": int, "filenames": list}
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                return {'file_count': 0, 'filenames': []}
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        files = [f for f in filelist if f['mimeType'] != 'application/vnd.google-apps.folder']
        return {'file_count': len(files), 'filenames': [f['title'] for f in files]}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'file_count': 0, 'filenames': []}

def get_chrome_font_size__4edc5823(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 0.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}

def get_gdrive_text_file__a652858db2e92fae77817389157c8edc(env, config: dict):
    """Get text file content from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'settings_file' and 'query_list'

    Returns:
        str: Content of the text file from Google Drive
    """
    from desktop_env.evaluators.getters.file import get_googledrive_file
    file_bytes = get_googledrive_file(env, config)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        print(f'Error decoding file: {e}')
        return ''

def get_table_position__b943f0fd(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_docx_last_table_dims__cbb3bae6(env, config: Dict[str, Any]):
    """Get dimensions of the last table in a docx file.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        dict: Dictionary with 'rows' and 'cols' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        if not doc.tables:
            return None
        last_table = doc.tables[-1]
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns)}
    except Exception as e:
        return None

def get_docx_table_data__d5cdfdb611de0d2892eeaddda9638ac9(env, config):
    """Extract table data from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists containing table cell values
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return []
        table = doc.tables[0]
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                try:
                    if '.' in cell_text:
                        row_data.append(float(cell_text))
                    else:
                        row_data.append(int(cell_text))
                except (ValueError, AttributeError):
                    row_data.append(cell_text)
            table_data.append(row_data)
        return table_data
    finally:
        os.unlink(tmp_path)

def get_table_position__d963ccf8(env, config):
    """Get the position and size of a table in a specific slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'slide_idx'

    Returns:
        dict: {'top': int, 'left': int, 'width': int, 'height': int} or None if no table found
    """
    vm_path = config.get('path')
    slide_idx = config.get('slide_idx', 0)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    import tempfile
    import os as os_module
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        tmp_file.write(file_bytes)
        tmp_path = tmp_file.name
    try:
        prs = Presentation(tmp_path)
        if slide_idx >= len(prs.slides):
            return None
        slide = prs.slides[slide_idx]
        for shape in slide.shapes:
            if shape.shape_type == 19:
                return {'top': shape.top, 'left': shape.left, 'width': shape.width, 'height': shape.height}
        return None
    finally:
        os_module.unlink(tmp_path)

def get_chrome_setting_value__eec6beed(env, config: Dict[str, str]):
    """
    Get the Translation settings setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['translate', 'enabled']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key)
                if setting_value is None:
                    logger.warning(f"[CHROME_SETTING] Key '{key}' not found in preferences, using Chrome default: True")
                    setting_value = True
                    break
            else:
                logger.warning(f"[CHROME_SETTING] Setting path incomplete at '{key}', using Chrome default: True")
                setting_value = True
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        raise Exception(f'Failed to read Chrome preferences file: {e}')

def get_googledrive_folder_file_count__1271f790(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get file information from a Google Drive folder including count and file types"""
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if len(folder_list) == 0:
            return {'count': 0, 'files': []}
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        files_info = []
        for file in file_list:
            files_info.append({'title': file.get('title', ''), 'mimeType': file.get('mimeType', ''), 'id': file.get('id', '')})
        return {'count': len(file_list), 'files': files_info}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'count': 0, 'files': []}

def get_extension_version__a366045b(env, config: Dict[str, Any]):
    """Extract the extension version from manifest.json.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension version, or None if not found
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    manifest_path = f'{extension_path}/manifest.json'
    try:
        content = env.controller.get_file(manifest_path)
        if content:
            manifest_data = json.loads(content)
            return manifest_data.get('version', None)
        return None
    except Exception as e:
        logger.error(f'Error reading extension version: {e}')
        return None

def get_docx_table_borders__5aa272b32299efe042e474c4fb5400ce(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get border information from tables in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        Dict with table border information: {'tables': [{'has_borders': bool}, ...]}
    """
    from docx import Document
    from docx.oxml.ns import qn
    file_path = config.get('path', '')
    if not file_path:
        return {'tables': []}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'tables': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        tables_info = []
        for table in doc.tables:
            has_borders = False
            tbl = table._element
            tblPr = tbl.tblPr
            if tblPr is not None:
                tblBorders = tblPr.find(qn('w:tblBorders'))
                if tblBorders is not None:
                    for border_name in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
                        border_elem = tblBorders.find(qn(f'w:{border_name}'))
                        if border_elem is not None:
                            val = border_elem.get(qn('w:val'))
                            if val and val != 'none':
                                has_borders = True
                                break
            tables_info.append({'has_borders': has_borders, 'rows': len(table.rows), 'columns': len(table.columns)})
        return {'tables': tables_info}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_chrome_setting_value__d84aae11(env, config: Dict[str, str]):
    """
    Get the Startup behavior setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['session', 'restore_on_startup']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, 1)
            else:
                setting_value = 1
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': 1}

def get_docx_table_content__0918fd6ddecb95fb671b970e39bbe53c(env, config: Dict[str, Any]) -> List[List[str]]:
    """Extract table content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        List of table rows, where each row is a list of cell values
    """
    file_path = config.get('path')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            all_tables.append(table_data)
        return all_tables
    finally:
        os.unlink(tmp_path)

def get_docx_table_data__0c283e204015cae1eab9b0792877891d(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract table data and document structure from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the VM file path

    Returns:
        Dictionary containing:
            - 'tables': List of tables, where each table is a list of rows
            - 'structure': List of document elements with their positions and types
                          (e.g., {'type': 'heading', 'text': 'Main Results', 'position': 5})
        Returns None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        tables_data = []
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cell_values = [cell.text.strip() for cell in row.cells]
                table_rows.append(cell_values)
            tables_data.append(table_rows)
        structure = []
        position = 0
        for element in doc.element.body:
            if element.tag.endswith('p'):
                para = Paragraph(element, doc)
                text = para.text.strip()
                if para.style.name.startswith('Heading') or 'Heading' in para.style.name:
                    structure.append({'type': 'heading', 'text': text, 'position': position, 'style': para.style.name})
                position += 1
            elif element.tag.endswith('tbl'):
                table_index = len([item for item in structure if item.get('type') == 'table'])
                structure.append({'type': 'table', 'table_index': table_index, 'position': position})
                position += 1
        return {'tables': tables_data, 'structure': structure}
    except Exception as e:
        print(f'Error reading docx: {e}')
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_docx_last_table_dims__c936c1ec(env, config: Dict[str, Any]):
    """Get dimensions of the last table and total table count in a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'rows', 'cols', 'table_count', and 'is_empty' keys, or None if no tables
    """
    filepath = config.get('path')
    if not filepath:
        return None
    try:
        doc = Document(filepath)
        table_count = len(doc.tables)
        if table_count == 0:
            return None
        last_table = doc.tables[-1]
        is_empty = True
        for row in last_table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    is_empty = False
                    break
            if not is_empty:
                break
        return {'rows': len(last_table.rows), 'cols': len(last_table.columns), 'table_count': table_count, 'is_empty': is_empty}
    except Exception as e:
        return None

def get_open_tab_count__a852a89d(env, config: dict):
    """Get number of open Chrome tabs.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used)

    Returns:
        Number of open tabs
    """
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    try:
        tabs = get_open_tabs_info(env, {})
        if isinstance(tabs, list):
            return len(tabs)
        return 0
    except Exception as e:
        logger.error(f'Failed to get tab count: {e}')
        return 0

def get_extension_enabled_state__6f71517e0c42749af1a6363d9f36e224(env, config: dict):
    """
    Check if a specific unpacked extension is enabled in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        dict: {'path': extension_path, 'enabled': True/False} or empty dict if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        target_path = config.get('extension_path', '')
        for (ext_id, ext_data) in all_extensions.items():
            if ext_data.get('path') == target_path:
                state = ext_data.get('state', 0)
                return {'path': target_path, 'enabled': state == 1}
        return {}
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return {}

def get_chrome_location_blocked__6591cb23(env, config: dict):
    """
    Check if Chrome's location access is blocked for all sites.
    This checks the 'profile.default_content_setting_values.geolocation' preference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used for this getter)

    Returns:
        str: "true" if location is blocked (value=2), "false" otherwise
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        geolocation_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('geolocation', 0)
        return 'true' if geolocation_setting == 2 else 'false'
    except Exception as e:
        logger.error(f'Error reading Chrome preferences: {e}')
        return 'false'

def get_googledrive_backup_check__c15ab8b3(env, config: Dict[str, Any]):
    """Comprehensive backup check.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: {"folder_exists": bool, "file_count": int, "expected_count": int}
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_path = config.get('folder_path', [])
        expected_count = config.get('expected_count', 0)
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and '{parent_id}' in parents and trashed = false"
            filelist = drive.ListFile({'q': query}).GetList()
            if len(filelist) == 0:
                return {'folder_exists': False, 'file_count': 0, 'expected_count': expected_count}
            parent_id = filelist[0]['id']
        query = f"'{parent_id}' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        eml_count = sum((1 for f in filelist if f.get('title', '').endswith('.eml')))
        return {'folder_exists': True, 'file_count': eml_count, 'expected_count': expected_count}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'folder_exists': False, 'file_count': 0, 'expected_count': 0}

def get_extension_description__db4e5321(env, config: Dict[str, Any]):
    """Extract the extension description from manifest.json.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' parameter

    Returns:
        str: Extension description, or None if not found
    """
    extension_path = config.get('extension_path', '/home/user/Desktop/helloExtension')
    manifest_path = f'{extension_path}/manifest.json'
    try:
        content = env.controller.get_file(manifest_path)
        if content:
            manifest_data = json.loads(content)
            return manifest_data.get('description', None)
        return None
    except Exception as e:
        logger.error(f'Error reading extension description: {e}')
        return None

def get_docx_arxiv_urls__82d38938(env, config):
    """Extract all arxiv.org URLs from a docx file.

    This function reads a .docx file and extracts all hyperlinks that contain
    'arxiv.org'. It looks through all paragraphs and runs in the document to
    find hyperlinked text.

    Args:
        env: Environment object with cache_dir and controller
        config: Configuration dict with:
            - path: Path to the docx file in the VM
            - dest (optional): Destination filename in cache

    Returns:
        list: List of arxiv.org URLs found in the document
    """
    import os
    import re
    from docx import Document
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    doc = Document(cache_path)
    arxiv_urls = []
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            element = run._element
            if element.tag.endswith('hyperlink'):
                r_id = element.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if r_id:
                    try:
                        rel = doc.part.rels[r_id]
                        url = rel.target_ref
                        if 'arxiv.org' in url:
                            arxiv_urls.append(url)
                    except:
                        pass
            parent = element.getparent()
            if parent is not None and parent.tag.endswith('hyperlink'):
                r_id = parent.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if r_id:
                    try:
                        rel = doc.part.rels[r_id]
                        url = rel.target_ref
                        if 'arxiv.org' in url:
                            arxiv_urls.append(url)
                    except:
                        pass
    for paragraph in doc.paragraphs:
        text = paragraph.text
        urls = re.findall('https?://[^\\s]*arxiv\\.org[^\\s]*', text)
        for url in urls:
            if url not in arxiv_urls:
                arxiv_urls.append(url)
    seen = set()
    unique_urls = []
    for url in arxiv_urls:
        if url not in seen:
            seen.add(url)
            unique_urls.append(url)
    return unique_urls

def get_extension_enabled_status__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get enabled status for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict mapping extension names to their enabled status
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        extensions_status = {}
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            name = ext_data.get('manifest', {}).get('name', '')
            state = ext_data.get('state', 0)
            enabled = state == 1
            extensions_status[name] = enabled
        return extensions_status
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file for extension status: {e}')
        return {}

def get_chrome_default_font_size__e289225d(env, config: dict):
    """Extract default font size from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        int: Default font size in pixels
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        default_font_size = data.get('webkit', {}).get('webprefs', {}).get('default_font_size', 16)
        logger.info(f'Chrome default font size: {default_font_size}')
        return int(default_font_size)
    except Exception as e:
        logger.error(f'Error getting Chrome font size: {e}')
        return 16

def get_default_web_browser__88e8e17a(env, config: dict):
    """Gets the default web browser on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'x-scheme-handler/http']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_html_file_and_chrome_tab__6920e35a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if HTML file exists and if it's opened in Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with file_exists (bool) and chrome_tabs (list of URLs)
    """
    html_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.html'
    result = env.controller.get_file(html_path)
    file_exists = False
    if result is not None and len(result) > 0:
        content = result.lower() if isinstance(result, str) else ''
        file_exists = ('<html' in content or '<!doctype html' in content or '<table' in content) and len(result) > 100
    if not file_exists and result is not None and (len(result) > 0):
        logger.warning(f'HTML file exists but may not be valid HTML or is too small (size: {len(result)} bytes)')
    elif file_exists:
        logger.info(f'Valid HTML file found at {html_path} (size: {len(result)} bytes)')
    from desktop_env.evaluators.getters.chrome import get_open_tabs_info
    tabs_info = get_open_tabs_info(env, {})
    chrome_tabs = []
    if tabs_info:
        chrome_tabs = [tab.get('url', '') for tab in tabs_info if isinstance(tab, dict)]
    return {'file_exists': file_exists, 'chrome_tabs': chrome_tabs, 'expected_file_path': html_path}

def get_chrome_setting_value__1be3beb2(env, config: Dict[str, str]):
    """
    Get the Download location setting from Chrome preferences.
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        setting_value = data
        for key in ['download', 'default_directory']:
            if isinstance(setting_value, dict):
                setting_value = setting_value.get(key, '/home/user/Downloads')
            else:
                setting_value = '/home/user/Downloads'
                break
        logger.info(f'[CHROME_SETTING] Retrieved setting value: {setting_value}')
        return {'setting_value': setting_value}
    except Exception as e:
        logger.error(f'Error getting Chrome setting: {e}')
        return {'setting_value': '/home/user/Downloads'}

def get_all_extension_paths__225a261e(env, config):
    """
    Get all loaded extension paths from Chrome Preferences file.

    This function reads the Chrome Preferences file and extracts the paths
    of all installed extensions (both from web store and unpacked).

    Args:
        env: Environment object with controller and vm_platform
        config: Configuration dict (not used in this function)

    Returns:
        list: List of extension paths as strings
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions_path = []
        all_extensions = data.get('extensions', {}).get('settings', {})
        for extension_id in all_extensions.keys():
            path = all_extensions[extension_id].get('path', '')
            if path:
                all_extensions_path.append(path)
        logger.info(f'Found {len(all_extensions_path)} extension paths')
        logger.debug(f'Extension paths: {all_extensions_path}')
        return all_extensions_path
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return []

def get_docx_table_info__fca84b62(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract table information from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: VM path to the DOCX file
            - dest: Local destination filename
            - table_index: Index of the table to check (default: -1 for last table)

    Returns:
        Dict with table information
    """
    vm_path = config.get('path', '/home/user/Desktop/Table_Of_Work_Effort_Instructions.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'error': 'Failed to retrieve file from VM'}
    dest = config.get('dest', 'table_doc.docx')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
    except Exception as e:
        return {'error': f'Failed to parse DOCX: {e}'}
    total_tables = len(doc.tables)
    table_index = config.get('table_index', -1)
    if table_index < 0:
        table_index = total_tables + table_index
    if table_index < 0 or table_index >= total_tables:
        return {'total_tables': total_tables, 'error': f'Table index {table_index} out of range'}
    table = doc.tables[table_index]
    rows = len(table.rows)
    columns = len(table.columns)
    cell_contents = []
    is_empty = True
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip()
            row_data.append(text)
            if text:
                is_empty = False
        cell_contents.append(row_data)
    return {'total_tables': total_tables, 'target_table': {'rows': rows, 'columns': columns, 'is_empty': is_empty, 'cell_contents': cell_contents}}

def get_chrome_javascript_setting__da00e3f2(env, config: Dict[str, str]):
    """
    Get the JavaScript content setting from Chrome preferences.

    Args:
        env: Desktop environment instance
        config: Configuration dictionary

    Returns:
        str: "allow" if JavaScript is enabled, "block" if disabled
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        javascript_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('javascript', 1)
        return 'block' if javascript_setting == 2 else 'allow'
    except Exception as e:
        logger.error(f'Error getting JavaScript setting: {e}')
        return 'allow'

def get_docx_table_data__b0edb7cf3cae7467b1751eb74a239c0d(env, config):
    """Extract table data and formatting from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict containing:
            - 'data': List of lists containing table cell values
            - 'formatting': List of lists containing formatting info for each cell
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'data': [], 'formatting': []}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.tables:
            return {'data': [], 'formatting': []}
        table = doc.tables[0]
        table_data = []
        table_formatting = []
        for row in table.rows:
            row_data = []
            row_formatting = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                try:
                    if '.' in cell_text:
                        row_data.append(float(cell_text))
                    else:
                        row_data.append(int(cell_text))
                except (ValueError, AttributeError):
                    row_data.append(cell_text)
                formatting_info = {'has_formatting': False, 'bold': False, 'italic': False, 'font_name': None, 'font_size': None}
                if cell.paragraphs:
                    for paragraph in cell.paragraphs:
                        if paragraph.runs:
                            for run in paragraph.runs:
                                if run.font:
                                    if run.font.bold:
                                        formatting_info['bold'] = True
                                        formatting_info['has_formatting'] = True
                                    if run.font.italic:
                                        formatting_info['italic'] = True
                                        formatting_info['has_formatting'] = True
                                    if run.font.name:
                                        formatting_info['font_name'] = run.font.name
                                        formatting_info['has_formatting'] = True
                                    if run.font.size:
                                        formatting_info['font_size'] = run.font.size.pt
                                        formatting_info['has_formatting'] = True
                row_formatting.append(formatting_info)
            table_data.append(row_data)
            table_formatting.append(row_formatting)
        return {'data': table_data, 'formatting': table_formatting}
    finally:
        os.unlink(tmp_path)

def get_writer_table_with_headers__cb84d8e55491bfd63123932189f4803b(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract table from Writer document including headers and verify data content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with headers, table data, and content verification metrics
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            if len(doc.tables) == 0:
                return {'error': 'No tables found'}
            table = doc.tables[0]
            all_rows = []
            non_empty_data_cells = 0
            total_data_cells = 0
            for row in table.rows:
                row_values = [cell.text.strip() for cell in row.cells]
                all_rows.append(row_values)
            data_rows = all_rows[1:] if len(all_rows) > 1 else []
            for row_values in data_rows:
                for cell_value in row_values:
                    total_data_cells += 1
                    if cell_value:
                        non_empty_data_cells += 1
            content_density = non_empty_data_cells / total_data_cells if total_data_cells > 0 else 0.0
            return {'headers': all_rows[0] if all_rows else [], 'data_rows': data_rows, 'total_rows': len(all_rows), 'non_empty_data_cells': non_empty_data_cells, 'total_data_cells': total_data_cells, 'content_density': content_density}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_extension_description__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get description for all installed extensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with 'extensions' (mapping extension names to descriptions) and 'all_names' (list of all extension names)
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    max_retries = 3
    retry_delay = 2
    for attempt in range(max_retries):
        try:
            content = env.controller.get_file(preference_file_path)
            data = json.loads(content)
            extensions_description = {}
            all_extensions = data.get('extensions', {}).get('settings', {})
            if not all_extensions and attempt < max_retries - 1:
                logger.warning(f'No extensions found in Preferences file, retrying in {retry_delay} seconds (attempt {attempt + 1}/{max_retries})')
                time.sleep(retry_delay)
                continue
            for ext_id in all_extensions.keys():
                ext_data = all_extensions[ext_id]
                name = ext_data.get('manifest', {}).get('name', '')
                description = ext_data.get('manifest', {}).get('description', '')
                if name:
                    extensions_description[name] = description
            all_extension_names = list(extensions_description.keys())
            if extensions_description:
                logger.info(f'Found {len(extensions_description)} extensions: {all_extension_names}')
                for (name, desc) in extensions_description.items():
                    logger.debug(f'  - {name}: {desc}')
            else:
                logger.warning('No extensions with valid names found in Preferences file')
            return {'extensions': extensions_description, 'all_names': all_extension_names}
        except Exception as e:
            if attempt < max_retries - 1:
                logger.warning(f'Error reading Chrome Preferences file (attempt {attempt + 1}/{max_retries}): {e}')
                time.sleep(retry_delay)
            else:
                logger.error(f'Error reading Chrome Preferences file for extension descriptions after {max_retries} attempts: {e}')
                return {'extensions': {}, 'all_names': []}
    return {'extensions': {}, 'all_names': []}

def get_table_count_only__e164fbd9(env, config):
    """Get first table structure and content to verify digraph table, and extract digraphs from document."""
    file_path = config.get('path')
    if not file_path:
        return {'table_count': 0, 'has_first_table': False, 'first_table_info': None, 'document_digraphs': []}
    file_data = env.controller.get_file(file_path)
    if not file_data:
        return {'table_count': 0, 'has_first_table': False, 'first_table_info': None, 'document_digraphs': []}
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        table_count = len(doc.tables)
        document_text = []
        first_table_found = False
        for element in doc.element.body:
            if element.tag.endswith('tbl'):
                if not first_table_found:
                    first_table_found = True
                    continue
                for row in element.xpath('.//w:tr', namespaces=element.nsmap):
                    for cell in row.xpath('.//w:tc', namespaces=element.nsmap):
                        cell_text = cell.xpath('string(.)', namespaces=element.nsmap).strip()
                        document_text.append(cell_text)
            elif element.tag.endswith('p'):
                para_text = element.xpath('string(.)', namespaces=element.nsmap).strip()
                document_text.append(para_text)
        document_digraphs = set()
        table_index = 0
        for element in doc.element.body:
            if element.tag.endswith('tbl'):
                if table_index == 0:
                    table_index += 1
                    continue
                for row in element.xpath('.//w:tr', namespaces=element.nsmap):
                    for cell in row.xpath('.//w:tc', namespaces=element.nsmap):
                        cell_text = cell.xpath('string(.)', namespaces=element.nsmap).strip()
                        if cell_text and len(cell_text) == 2 and cell_text.isalpha():
                            document_digraphs.add(cell_text.lower())
                        quoted_pattern = '["\\\']([a-zA-Z]{2})["\\\']|^([a-zA-Z]{2})$|\\b([a-zA-Z]{2})\\b'
                        matches = re.findall(quoted_pattern, cell_text)
                        for match_tuple in matches:
                            for match in match_tuple:
                                if match and len(match) == 2 and match.isalpha():
                                    document_digraphs.add(match.lower())
                table_index += 1
        full_text = ' '.join(document_text)
        explicit_patterns = ['["\\\']([a-zA-Z]{2})["\\\']', '\\(([a-zA-Z]{2})\\)', ':\\s*([a-zA-Z]{2})\\b', '-\\s*([a-zA-Z]{2})\\b', '^\\s*([a-zA-Z]{2})\\s*$']
        for pattern in explicit_patterns:
            matches = re.findall(pattern, full_text, re.MULTILINE | re.IGNORECASE)
            for match in matches:
                if len(match) == 2 and match.isalpha():
                    document_digraphs.add(match.lower())
        document_digraphs_list = sorted(list(document_digraphs))
        logger.info(f'Found {len(document_digraphs_list)} digraphs in document: {document_digraphs_list[:20]}')
        first_table_info = None
        if table_count > 0:
            first_table = doc.tables[0]
            is_at_beginning = True
            for element in doc.element.body:
                if element.tag.endswith('tbl'):
                    break
                elif element.tag.endswith('p'):
                    para_text = element.xpath('string(.)').strip()
                    if len(para_text) > 10:
                        is_at_beginning = False
                        break
            row_count = len(first_table.rows)
            column_count = len(first_table.columns) if first_table.rows else 0
            cells = []
            if row_count > 0:
                first_row = first_table.rows[0]
                cells = [cell.text.strip() for cell in first_row.cells]
            are_digraphs = all((len(cell) == 2 and cell.isalpha() for cell in cells if cell))
            first_table_info = {'position': 'first' if is_at_beginning else 'not_first', 'row_count': row_count, 'column_count': column_count, 'cells': cells, 'are_digraphs': are_digraphs, 'is_single_row': row_count == 1}
        os.unlink(tmp_path)
        logger.info(f'Found {table_count} tables, first table: {first_table_info}')
        return {'table_count': table_count, 'has_first_table': table_count > 0, 'first_table_info': first_table_info, 'document_digraphs': document_digraphs_list}
    except Exception as e:
        logger.error(f'Error: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return {'table_count': 0, 'has_first_table': False, 'first_table_info': None, 'document_digraphs': []}

def get_chrome_sansserif_font__2cf2a146(env, config: Dict[str, str]):
    """Get Chrome sans-serif font family setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with sansserif_font_family value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        fonts = data.get('webkit', {}).get('webprefs', {}).get('fonts', {}).get('sansserif', {})
        sansserif_font = fonts.get('Zyyy', 'Arial')
        return {'sansserif_font_family': sansserif_font}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'sansserif_font_family': 'Arial'}

def get_docx_table_content__28c7c94bbd8c2d75a3d0fb78ff321239(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract table content from a DOCX file with section information.

    This getter reads a DOCX file and extracts all tables with their position
    relative to document sections/headings.

    Args:
        env: DesktopEnv instance with controller
        config: Configuration dict with 'path' key pointing to the DOCX file on VM

    Returns:
        Dictionary with tables data, or None if error occurs
        Format: {
            'num_tables': int,
            'tables': [
                {
                    'num_rows': int,
                    'num_cols': int,
                    'data': [[cell_text, ...], ...],
                    'preceding_heading': str or None,
                    'section_context': str or None
                },
                ...
            ]
        }
    """
    try:
        file_path = config.get('path', '')
        if not file_path:
            logger.error('No path specified in config')
            return None
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.error(f'Failed to get file from VM: {file_path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            table_positions = {}
            for (i, table) in enumerate(doc.tables):
                for (j, element) in enumerate(doc.element.body):
                    if element == table._element:
                        table_positions[id(table)] = j
                        break
            tables_data = []
            for table in doc.tables:
                table_info = {'num_rows': len(table.rows), 'num_cols': len(table.columns) if table.rows else 0, 'data': [], 'preceding_heading': None, 'section_context': None}
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_info['data'].append(row_data)
                table_pos = table_positions.get(id(table), -1)
                if table_pos >= 0:
                    for i in range(table_pos - 1, -1, -1):
                        element = doc.element.body[i]
                        if element.tag.endswith('p'):
                            para_idx = None
                            for (idx, p) in enumerate(doc.paragraphs):
                                if p._element == element:
                                    para_idx = idx
                                    break
                            if para_idx is not None:
                                para = doc.paragraphs[para_idx]
                                if para.style and para.style.name and ('Heading' in para.style.name):
                                    table_info['preceding_heading'] = para.text.strip()
                                    table_info['section_context'] = para.text.strip()
                                    break
                                elif para.text.strip() and len(para.runs) > 0:
                                    if all((run.bold for run in para.runs if run.text.strip())):
                                        if not table_info['preceding_heading']:
                                            table_info['preceding_heading'] = para.text.strip()
                                            table_info['section_context'] = para.text.strip()
                tables_data.append(table_info)
            result = {'num_tables': len(tables_data), 'tables': tables_data}
            return result
        finally:
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting DOCX table content: {e}')
        return None

def get_chrome_experiments_contains__0a40109ea287ab7bbd8cd9175a7ce6a5(env, config: Dict[str, str]):
    """
    Get enabled Chrome experiments and return them as a list.
    This getter is used to check if specific experiments are enabled.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, but required by framework)

    Returns:
        List of enabled experiment names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Local State'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enabled_labs_experiments = data.get('browser', {}).get('enabled_labs_experiments', [])
        experiment_names = [exp.split('@')[0] for exp in enabled_labs_experiments]
        return experiment_names
    except Exception as e:
        logger.error(f'Error getting enabled experiments: {e}')
        return []

def get_extension_folder_exists__ae6416e4(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if extension folder exists AND if extension is loaded into Chrome.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' key

    Returns:
        Dict with:
            - 'folder_exists': boolean - Whether the folder exists
            - 'folder_path': str - Path to the folder
            - 'extension_installed': boolean - Whether extension is loaded in Chrome
            - 'installed_path': str - Path in Chrome Preferences (if installed)
    """
    folder_path = config.get('folder_path', '')
    result = {'folder_exists': False, 'folder_path': folder_path, 'extension_installed': False, 'installed_path': ''}
    if not folder_path:
        return result
    try:
        folder_check = env.controller.run_bash_script(f'[ -d "{folder_path}" ] && echo "exists" || echo "not_exists"', timeout=10)
        output = folder_check.get('output', '').strip()
        result['folder_exists'] = output == 'exists'
        logger.info(f"Extension folder exists: {result['folder_exists']}")
    except Exception as e:
        logger.error(f'Error checking extension folder: {e}')
        result['folder_exists'] = False
    os_type = env.vm_platform
    try:
        if os_type == 'Windows':
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
        elif os_type == 'Darwin':
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
        elif os_type == 'Linux':
            if 'arm' in platform.machine():
                preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
            else:
                preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
        else:
            logger.error(f'Unsupported operating system: {os_type}')
            return result
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for (extension_id, extension_data) in all_extensions.items():
            extension_path = extension_data.get('path', '')
            if extension_path and folder_path in extension_path:
                result['extension_installed'] = True
                result['installed_path'] = extension_path
                logger.info(f'Extension found in Chrome Preferences: {extension_id} at {extension_path}')
                break
        if not result['extension_installed']:
            logger.warning(f'Extension folder exists but is NOT loaded in Chrome')
            logger.info(f'Checked {len(all_extensions)} installed extensions')
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        result['extension_installed'] = False
    return result

def get_gdrive_file_check__db0d2d11(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in a Google Drive folder.

    Args:
        env: Desktop environment instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - folder_query: Query to find the folder
            - file_query: Query to find the file

    Returns:
        Dict with file existence and folder information
    """
    sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', '..', '..'))
    from desktop_env.evaluators.google_drive import GoogleDriveEvaluator
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_query = config.get('folder_query', '')
    file_query = config.get('file_query', '')
    try:
        gd_evaluator = GoogleDriveEvaluator(settings_file)
        folder_results = gd_evaluator.search_file(folder_query)
        folder_exists = len(folder_results) > 0
        folder_id = folder_results[0]['id'] if folder_exists else None
        file_results = gd_evaluator.search_file(file_query)
        file_exists = len(file_results) > 0
        in_folder = False
        if file_exists and folder_exists:
            for file_item in file_results:
                parents = file_item.get('parents', [])
                if folder_id in parents:
                    in_folder = True
                    break
        result = {'file_exists': file_exists, 'folder_exists': folder_exists, 'in_folder': in_folder, 'folder_id': folder_id, 'file_count': len(file_results)}
        logger.info(f'Google Drive check result: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking Google Drive: {e}')
        return {'file_exists': False, 'folder_exists': False, 'in_folder': False, 'folder_id': None, 'file_count': 0}

def get_chrome_serif_font__ddb77dce(env, config: Dict[str, str]):
    """Get Chrome serif font family setting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Dict with serif_font_family value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        fonts = data.get('webkit', {}).get('webprefs', {}).get('fonts', {}).get('serif', {})
        serif_font = fonts.get('Zyyy', 'Times New Roman')
        return {'serif_font_family': serif_font}
    except Exception as e:
        logger.error(f'Error: {e}')
        return {'serif_font_family': 'Times New Roman'}

def get_webext_manifest__dea4a5ea3d588b414bf151fee72b35b0(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract manifest.json from a web extension project directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to manifest.json on VM)

    Returns:
        Dict containing the manifest JSON, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get manifest file from VM: {path}')
            return None
        manifest = json.loads(file_bytes.decode('utf-8'))
        logger.info(f'Successfully loaded manifest from {path}')
        return manifest
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON from {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading manifest from {path}: {e}')
        return None

def get_chrome_experiments_multi__e998f78abb27064086318477b860256b(env, config: Dict[str, str]):
    """
    Get enabled Chrome experiments and return them as a list.
    This getter is used to check if multiple specific experiments are enabled.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, but required by framework)

    Returns:
        List of enabled experiment names
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Local State'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Local State'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Local State'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Local State'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        enabled_labs_experiments = data.get('browser', {}).get('enabled_labs_experiments', [])
        experiment_names = [exp.split('@')[0] for exp in enabled_labs_experiments]
        return experiment_names
    except Exception as e:
        logger.error(f'Error getting enabled experiments: {e}')
        return []

def get_chrome_fixed_font_size__d608837f75cec39dc6023178df898af7(env, config: Dict[str, str]):
    """
    Get the fixed-width (monospace) font size setting from Chrome preferences.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Dictionary containing default_fixed_font_size value
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                                'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        webprefs = data.get('webkit', {}).get('webprefs', {})
        default_fixed_font_size = webprefs.get('default_fixed_font_size', 13)
        return {'default_fixed_font_size': default_fixed_font_size}
    except Exception as e:
        logger.error(f'Error getting fixed font size: {e}')
        return {'default_fixed_font_size': 13}

def get_webext_dir__b9b28e4a(env, config):
    """
    Get the extension directory path on VM and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        str: Directory path where the extension is located
    """
    dir_path = config.get('path', '/home/user/Projects')
    ext_name = config.get('ext_name', 'my-extension')
    full_path = f'{dir_path}/{ext_name}'
    return full_path

def get_sorted_table_data__857e1276(env, config: Dict[str, Any]) -> List[List[Any]]:
    """Extract sorted table data from a specific range in Excel file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path', 'sheet', 'start_row', 'start_col', 'num_rows', 'num_cols'

    Returns:
        List of lists containing cell values from the specified range
    """
    path = config.get('path', '/home/user/Students_Class_Subject_Marks.xlsx')
    sheet_idx = config.get('sheet', 0)
    start_row = config.get('start_row', 2)
    start_col = config.get('start_col', 2)
    num_rows = config.get('num_rows', 4)
    num_cols = config.get('num_cols', 5)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        if isinstance(sheet_idx, int):
            ws = wb.worksheets[sheet_idx]
        else:
            ws = wb[sheet_idx]
        result = []
        for row in range(start_row, start_row + num_rows):
            row_data = []
            for col in range(start_col, start_col + num_cols):
                cell_value = ws.cell(row, col).value
                row_data.append(cell_value)
            result.append(row_data)
        wb.close()
        return result
    finally:
        os.unlink(tmp_path)

def get_recreation_devilsgarden_html__fa1e76c31141f93d38de11c4bb8239cf(env, config: Dict[str, Any]):
    """
    Get HTML content from recreation.gov page for Devil's Garden search.
    This getter verifies:
    1. The page/URL contains "Devil's Garden" reference
    2. The reservation table is present
    3. Reservation data exists
    """
    logger.info(f"[RECREATION_DEVILSGARDEN] Starting recreation.gov page processing for Devil's Garden")
    logger.debug(f'[RECREATION_DEVILSGARDEN] Config: {config}')
    host = env.vm_ip
    port = env.chromium_port
    server_port = env.server_port
    use_proxy = env.current_use_proxy
    remote_debugging_url = f'http://{host}:{port}'
    backend_url = f'http://{host}:{server_port}'
    max_retries = 3
    timeout_ms = 60000
    for attempt in range(max_retries):
        try:
            logger.info(f'[RECREATION_DEVILSGARDEN] Attempt {attempt + 1}/{max_retries}')
            with sync_playwright() as p:
                try:
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_DEVILSGARDEN] Successfully connected to existing Chrome instance')
                except Exception as e:
                    logger.warning(f'[RECREATION_DEVILSGARDEN] Failed to connect to existing Chrome instance: {e}')
                    logger.info(f'[RECREATION_DEVILSGARDEN] Starting new Chrome instance...')
                    app = 'chromium' if 'arm' in platform.machine() else 'google-chrome'
                    command = [app, '--remote-debugging-port=1337', '--no-sandbox']
                    if use_proxy:
                        command.append(f'--proxy-server=127.0.0.1:18888')
                    logger.info(f"[RECREATION_DEVILSGARDEN] Starting browser with command: {' '.join(command)}")
                    payload = json.dumps({'command': command, 'shell': False})
                    headers = {'Content-Type': 'application/json'}
                    requests.post(backend_url + '/setup' + '/launch', headers=headers, data=payload)
                    time.sleep(5)
                    browser = p.chromium.connect_over_cdp(remote_debugging_url)
                    logger.info(f'[RECREATION_DEVILSGARDEN] Successfully connected to new Chrome instance')
                if len(browser.contexts) == 0 or len(browser.contexts[0].pages) == 0:
                    logger.error(f'[RECREATION_DEVILSGARDEN] No active pages found')
                    return None
                page = browser.contexts[0].pages[0]
                current_url = page.url
                logger.info(f'[RECREATION_DEVILSGARDEN] Current URL: {current_url}')
                content = page.content()
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'html.parser')
                result = {'location_verified': False, 'reservation_table_present': False, 'url': current_url, 'page_title': '', 'has_availability_data': False, 'dates_sorted': False, 'earliest_reservation_identified': False, 'reservation_dates': []}
                page_text = soup.get_text().lower()
                url_lower = current_url.lower()
                if 'devil' in url_lower or 'devils' in url_lower or '232449' in url_lower:
                    result['location_verified'] = True
                    logger.info(f'[RECREATION_DEVILSGARDEN] Location verified in URL')
                elif "devil's garden" in page_text or 'devils garden' in page_text or 'devil garden' in page_text:
                    result['location_verified'] = True
                    logger.info(f'[RECREATION_DEVILSGARDEN] Location verified in page content')
                else:
                    logger.warning(f"[RECREATION_DEVILSGARDEN] Could not verify Devil's Garden location")
                title_tag = soup.find('title')
                if title_tag:
                    result['page_title'] = title_tag.get_text().strip()
                    logger.info(f"[RECREATION_DEVILSGARDEN] Page title: {result['page_title']}")
                selector = config.get('selector', 'class')
                class_name = config.get('class', 'camp-sortable-column-header')
                if selector == 'class':
                    elements = soup.find_all(class_=class_name)
                    logger.info(f"[RECREATION_DEVILSGARDEN] Found {len(elements)} elements with class '{class_name}'")
                    if len(elements) >= 2:
                        result['reservation_table_present'] = True
                        logger.info(f'[RECREATION_DEVILSGARDEN] Reservation table confirmed with {len(elements)} headers')
                availability_indicators = [soup.find_all(class_='availability-status'), soup.find_all(class_='campsite-row'), soup.find_all(class_='rec-availability'), soup.find_all('td', class_=lambda x: x and 'available' in x.lower() if x else False), soup.find_all('div', class_=lambda x: x and 'reservation' in x.lower() if x else False)]
                for indicator_list in availability_indicators:
                    if indicator_list and len(indicator_list) > 0:
                        result['has_availability_data'] = True
                        logger.info(f'[RECREATION_DEVILSGARDEN] Found availability data: {len(indicator_list)} elements')
                        break
                if not result['has_availability_data']:
                    date_elements = soup.find_all(class_=lambda x: x and 'date' in x.lower() if x else False)
                    if len(date_elements) >= 3:
                        result['has_availability_data'] = True
                        logger.info(f'[RECREATION_DEVILSGARDEN] Found date elements indicating availability view')
                from datetime import datetime
                import re
                reservation_dates = []
                date_patterns = ['\\b(\\d{1,2})/(\\d{1,2})/(\\d{2,4})\\b', '\\b(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\\s+(\\d{1,2}),?\\s+(\\d{4})\\b', '\\b(\\d{4})-(\\d{2})-(\\d{2})\\b']
                date_containers = soup.find_all(['td', 'th', 'div', 'span'], class_=lambda x: x and any((keyword in str(x).lower() for keyword in ['date', 'day', 'calendar', 'availability'])) if x else False)
                table_rows = soup.find_all('tr')
                date_containers.extend(table_rows)
                for element in date_containers:
                    text = element.get_text()
                    for pattern in date_patterns:
                        matches = re.findall(pattern, text)
                        for match in matches:
                            try:
                                if len(match) == 3:
                                    if '/' in text:
                                        (month, day, year) = (int(match[0]), int(match[1]), int(match[2]))
                                        if year < 100:
                                            year += 2000
                                        date_obj = datetime(year, month, day)
                                    elif '-' in text:
                                        (year, month, day) = (int(match[0]), int(match[1]), int(match[2]))
                                        date_obj = datetime(year, month, day)
                                    else:
                                        months = {'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4, 'may': 5, 'jun': 6, 'jul': 7, 'aug': 8, 'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12}
                                        month = months[match[0].lower()[:3]]
                                        day = int(match[1])
                                        year = int(match[2])
                                        date_obj = datetime(year, month, day)
                                    reservation_dates.append(date_obj)
                                    logger.info(f"[RECREATION_DEVILSGARDEN] Found reservation date: {date_obj.strftime('%Y-%m-%d')}")
                            except (ValueError, KeyError) as e:
                                continue
                reservation_dates = sorted(list(set(reservation_dates)))
                result['reservation_dates'] = [d.strftime('%Y-%m-%d') for d in reservation_dates]
                logger.info(f'[RECREATION_DEVILSGARDEN] Extracted {len(reservation_dates)} unique reservation dates')
                if len(reservation_dates) >= 2:
                    is_sorted = all((reservation_dates[i] <= reservation_dates[i + 1] for i in range(len(reservation_dates) - 1)))
                    if is_sorted:
                        result['dates_sorted'] = True
                        logger.info(f'[RECREATION_DEVILSGARDEN] Dates are sorted in ascending order (earliest first)')
                    else:
                        logger.warning(f'[RECREATION_DEVILSGARDEN] Dates are NOT sorted in ascending order')
                if len(reservation_dates) > 0:
                    earliest_date_str = reservation_dates[0].strftime('%Y-%m-%d')
                    highlighted_indicators = [soup.find_all(class_=lambda x: x and any((keyword in str(x).lower() for keyword in ['selected', 'active', 'highlight', 'focus'])) if x else False), soup.find_all(attrs={'aria-selected': 'true'}), soup.find_all(attrs={'data-selected': 'true'})]
                    for indicator_list in highlighted_indicators:
                        for element in indicator_list:
                            element_text = element.get_text()
                            if earliest_date_str in element_text or reservation_dates[0].strftime('%m/%d/%Y') in element_text:
                                result['earliest_reservation_identified'] = True
                                logger.info(f'[RECREATION_DEVILSGARDEN] Earliest reservation ({earliest_date_str}) is highlighted/selected')
                                break
                        if result['earliest_reservation_identified']:
                            break
                    if not result['earliest_reservation_identified'] and result['dates_sorted'] and result['reservation_table_present']:
                        result['earliest_reservation_identified'] = True
                        logger.info(f'[RECREATION_DEVILSGARDEN] Earliest reservation identifiable at top of sorted table')
                logger.info(f'[RECREATION_DEVILSGARDEN] Final result: {result}')
                return result
        except Exception as e:
            logger.error(f'[RECREATION_DEVILSGARDEN] Attempt {attempt + 1} failed: {e}')
            if attempt < max_retries - 1:
                logger.info(f'[RECREATION_DEVILSGARDEN] Retrying in 2 seconds...')
                time.sleep(2)
            else:
                logger.error(f'[RECREATION_DEVILSGARDEN] All retries exhausted')
                return None
    return None

def get_file_extension__e13be972(env, config: dict):
    """Get the file extension of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: File extension (e.g., 'xlsx', 'pdf')
    """
    path = config.get('path', '/home/user/Desktop/invoice.xlsx')
    (_, ext) = os.path.splitext(path)
    ext = ext.lstrip('.')
    result = env.controller.run_bash_script(f'test -f "{path}" && echo "exists" || echo "not_found"', timeout=10)
    output = result.get('output', '').strip()
    if output == 'exists':
        return ext
    else:
        return ''

def get_googledrive_folder_info__a18b8359(env, config: Dict[str, Any]):
    """Get information about a Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: path to Google Drive settings
            - folder_name: name of folder to check

    Returns:
        dict: {"exists": bool, "file_count": int}
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_name = config.get('folder_name', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        query = f"title = '{folder_name}' and mimeType = 'application/vnd.google-apps.folder' and 'root' in parents and trashed = false"
        filelist = drive.ListFile({'q': query}).GetList()
        if len(filelist) == 0:
            return {'exists': False, 'file_count': 0}
        folder_id = filelist[0]['id']
        file_query = f"'{folder_id}' in parents and trashed = false"
        file_list = drive.ListFile({'q': file_query}).GetList()
        file_count = sum((1 for f in file_list if f['mimeType'] != 'application/vnd.google-apps.folder'))
        return {'exists': True, 'file_count': file_count}
    except Exception as e:
        logger.error(f'Error checking Google Drive folder: {e}')
        return {'exists': False, 'file_count': 0}

def get_extension_version__407be0458b7b234fb5401d66a10f5221(env, config):
    """Get version information for a specific extension by name.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_name' key

    Returns:
        str: Version number of the extension, or empty string if not found
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    target_name = config.get('extension_name', '')
    if not target_name:
        return ''
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for ext_id in all_extensions.keys():
            ext_data = all_extensions[ext_id]
            if 'manifest' in ext_data:
                name = ext_data['manifest'].get('name', '')
                if name.lower() == target_name.lower():
                    version = ext_data['manifest'].get('version', '')
                    return version
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return ''

def get_gdrive_file_list__3dcb78c90ac64690f9f399090d59db08(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a Google Drive folder by path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_path: List of folder names forming the path (e.g., ['mail_archive'])

    Returns:
        List of filenames found in the specified folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_path = config.get('folder_path', [])
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        parent_id = 'root'
        for folder_name in folder_path:
            query = f"'{parent_id}' in parents and title='{folder_name}' and mimeType='application/vnd.google-apps.folder' and trashed=false"
            file_list = drive.ListFile({'q': query}).GetList()
            if not file_list:
                logger.warning(f"Folder '{folder_name}' not found in path {folder_path}")
                return []
            parent_id = file_list[0]['id']
        query = f"'{parent_id}' in parents and trashed=false"
        file_list = drive.ListFile({'q': query}).GetList()
        filenames = [f['title'] for f in file_list]
        logger.info(f'Found {len(filenames)} files in Google Drive folder: {filenames}')
        return sorted(filenames)
    except Exception as e:
        logger.error(f'Error accessing Google Drive: {e}')
        return []

def get_chrome_font_size__9b3236ac(env, config):
    """
    Get Chrome's default font size from Preferences.
    This is a variation-specific getter for task variation 7.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)
    
    Returns:
        dict: Font size settings containing default_font_size, default_fixed_font_size, minimum_font_size
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'), 'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        search_engine = data.get('webkit', {}).get('webprefs', {'default_fixed_font_size': 13, 'default_font_size': 16, 'minimum_font_size': 13})
        logger.info(f'Retrieved Chrome font settings: {search_engine}')
        return search_engine
    except Exception as e:
        logger.error(f'Error retrieving Chrome font size: {e}')
        return {'default_fixed_font_size': 13, 'default_font_size': 16}
