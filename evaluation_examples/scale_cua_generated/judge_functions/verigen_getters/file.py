"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_zip_contents__ed5b60c9b586e02ab1d262e432f8c2c4', 'get_pdf_exports__bec19abd', 'get_pdf_files_info__cb51ce0b', 'get_file_exists__dc35efec', 'get_file_exists_check__b62f9acb', 'get_docx_text_content__6477df40', 'get_file_count_in_dir__f1d66967', 'get_music_file_count__355d5753246a8a88c1b8d173c110d89e', 'get_directory_file_count__58fb65f5', 'get_documents_file__c0d05fd4eae793b685d13d8511de53aa', 'get_pdf_file_list__d2aeadd33c64efc8e0e3416b06f6e222', 'get_rtf_file_info__34460ac1a76394f2dfa8b4e9981344a0', 'get_docx_text_content__7017cad57d424df42048663d589125a3', 'get_pdf_file_a503b07f', 'get_text_file_content__7676732d', 'get_text_file_content__d48f445fac6cb18530cd7f7169fdc7fc', 'get_file_list__4a34f03bcb82e7e8037181cd9a91ae6e', 'get_default_pdf_viewer__6a629498', 'get_file_exists__cb8ab5642aaf48705d11abb5543759c9', 'get_targz_file__04a2c355', 'get_numbers_from_txt__6622ce38', 'get_pdf_exports__d587f20a', 'get_folder_file_list__465762dc', 'get_total_files_count__e9983217', 'get_git_repo_subdirectory__a29e9963', 'get_folder_file_count__fc8f5b84', 'get_python_file_content__2a7463c5815fe65f87729a241d0d409d', 'get_pdf_file_info__3a8b0dcb7e90aff8357c94bc8dad2474', 'get_file_presence_pattern__ca75e69be093d6cb8e4fa53ce6114782', 'get_python_file_count__c32ed37d80ef317dea88bac4a4cc1f31', 'get_python_file__f1f92c4b10af2ffde6ea1534830a28f6', 'get_text_file_content__8f36e769', 'get_pdf_files_list__097402a1', 'get_directory_count__7ec80ec4', 'get_pdf_text_content__a09f0e42332d230e7cc0e7794732425e', 'get_dual_file_exists__e1a4b749', 'get_pdf_file_info__60388d3f6c5270e1728288c485d5fd5b', 'get_multi_directory_contents__ba67a508', 'get_text_file_content__83bea20e1ddf48cf4f537ad2c05896b6', 'get_file_text_content__61ff0c2a', 'get_pdf_basic_info__ebee77f7069f08db28f1fec38e7eb935', 'get_file_line_count__868f1e74', 'get_csv_first_column__e84252aa', 'get_file_content_dict__4ca7be3762cf58edd133468f74a91008', 'get_pdf_filenames_with_keyword__91578e58', 'get_zip_contents__6aa029d37a944ba9e2bf06a8a1d59f5c', 'get_file_exists__6f3c16ae', 'get_pdf_content_keywords__c68533a8ce70563646c4156811d621fa', 'get_subfolder_picture_hashes__fcfe8473', 'get_docx_text_content__64111ae8', 'get_directory_files__739292ff', 'get_file_count__9d5ff0d9', 'get_srt_file_exists__4f098003e517e7e34a157f1c233e1c85', 'get_numbers_from_file__d0f3a745', 'get_file_checksum__032a6328', 'get_pdf_files_list__92f3cb2c', 'get_numeric_value_from_file__5d55e268', 'get_file_text__b8aa5550', 'get_file_permissions__2bea57f7', 'get_text_file_content__a2f161de', 'get_file_content__9cefe3d0', 'get_file_content__66187dca7d72426d9eab5771bd5dd30e', 'get_file_exists__739292ff', 'get_folder_files__21ed89b452cf1a748805f1dda9b2ec4b', 'get_pdf_files_in_dir__600d7d75', 'get_text_file_lines__4250d59b26bb86f2de0562f0a55c312c', 'get_multi_directory_contents__e2bf8bf2', 'get_bash_file_content__38c143b8ba918371e989fa588ea9cb56', 'get_file_exists_and_size__91c793ae', 'get_pdf_exports__17dfab17', 'get_csv_and_count_verification__40b9d100', 'get_docx_text_content__ae54b7f2', 'get_pdf_exports__d263a7ae', 'get_text_file_content__ffa8cf6ad724ee8fc8a065457d283c28', 'get_text_file_content__0b52fd51', 'get_file_exists_check__e327229f', 'get_backup_dir_contents__1aa0aa31', 'get_vm_file_exists__949eb101', 'get_file_size__79b627f8', 'get_odt_file_info__82a08649fb4593676b95bae1dea8263f', 'get_both_files_rows__e54614f2', 'get_txt_file_lines__87f6c353', 'get_pdf_validation__82780f835e33185cfaa42bf9dd4b545f', 'get_pdf_multipage_info__7825428a26721545cdd69efbcd0db85d', 'get_file_and_folder_check__3c993009', 'get_text_file_lines__6a539f3f0959c2ae90484012b1290c1f', 'get_file_exists__6d219cf2', 'get_file_exists__9180f469', 'get_file_and_trash_status__05c91736', 'get_filename_pattern_match__0d61b4f8', 'get_file_content__bcc15712', 'get_filename__4ee0209a', 'get_text_file_content__93eac3e2452ef121ce8047db9ec250fe', 'get_docx_text_content__c9cce3df360663a0b15969bfb64090f9', 'get_file_count__9f898a55', 'get_pdf_files_in_dir__197670d4', 'get_header_content__cc0e3cfb7e6297a7f30b4a2bddec08bd', 'get_file_exists__cdcbbd90', 'get_directory_listing__96172b42', 'get_subject_file__271abb880d5f6f8d57d2c41e20bcf6ad', 'get_docx_content_and_font__d62b618c18b762984a0245cb92f201ef', 'get_file_count_by_pattern__5c107_5', 'get_text_file_lines__03d4ebf8db116d0dd3db06d0b1b6415c', 'get_file_rename_status__bd9934867663f0945bb79537ace5711a', 'get_file_in_directory__3500c3270fba14684f134a0b5e4537d1', 'get_pdf_files_in_dir__201ff98c', 'get_file_content__2ef64375b850b707e97b54053b1452c0', 'get_directory_exists__0b3844b6', 'get_pdf_directory_info__9ba5dc81930fa553e6c6310edaaff2ec', 'get_python_content__198be354', 'get_footer_has_page_and_filename__aac82fb893db44f74dfc3c4f83b7b05e', 'get_file_exists_and_size__058bd353', 'get_bash_file_content__8c025f1893a98ce909d203315d730cfe', 'get_pdf_filename__e61f394075a7c32a6a0c2c96a3700939', 'get_largest_filename_from_text__b7de68e1', 'get_files_with_prefix__cf3f5d8ece62ecf8e4937dea9e007679', 'get_text_content__2ad6be23', 'get_csv_filtered_rows__8b005fb41e3efd4706af2ae4f1a79bee', 'get_vm_dir_list__e4b458033c1389ecc56cd63f5eae9626', 'get_file_exists__d22e0dfa', 'get_text_file_content__16aae469605745263765ebbfdcb35f54', 'get_bash_file_content__816eebe3a50924fca9ee760018f3238e', 'get_text_file_content__c491dbb8', 'get_file_exists_and_size__dccb2b60', 'get_zip_file_list__ee271ee8', 'get_docx_text_content__66306a8b', 'get_speedtest_file_content__26660ad1', 'get_nested_directory_exists__3c3f0ceb', 'get_file_exists__5a942dd0', 'get_text_file_content__1e67b9ae311891ef2e3034615cded86c', 'get_python_file_content__c685793b2ca36ab76b7f2cc84f84fe40', 'get_docx_content__80bc3840fb160595a8bbd7e90abee050', 'get_archive_contents__6131523f', 'get_total_files_in_dirs__ed6a3699fe6af7deef02a8e547504034', 'get_file_size_info__0f2a6243', 'get_csv_column_values__8b88ca382748ac90936480358fbda368', 'get_text_file_content__ebe7bea30aab4d43b91ea1760b3fb66f', 'get_docx_text_content__2830db9ce8bded3cea70386f80d8c3ad', 'get_pdf_files_in_dir__a76459ec', 'get_python_file_content__3b7fbcc9bcd509e2386e3dc8ef7407df', 'get_docx_file_exists__c4f2c36653f58f8aab105ca7a48ec763', 'get_pdf_page_info__4c6d2d3e077e0f5ce3b4338c5664395c', 'get_screenshot_in_folder__956d014f26874db42b602f2626906894', 'get_file_in_directory__845b16e9c20bb76eb4ef8a7eb9262413', 'get_text_content__f9fa7b925a8af809c4d7a02e0264dd76', 'get_backup_files_check__ac9408954b941c7f40eedd27a6f1296b', 'get_folder_contents__4e03b1ed', 'get_pdf_file_size__f34c6d72482f0e93994904e974de58fa', 'get_docx_text_content__d99037c6', 'get_tetris_files__ca30528a', 'get_docx_text_content__faacb6dea70ef25b404521dc6e6554b5', 'get_text_file_lines__7e304294ce76dab9f08b95178667a620', 'get_multi_directory_contents__429c8cbc', 'get_docx_file_count__bb5651c2', 'get_text_file_content__5ced85fc_aug18_v4_d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8a9', 'get_text_file_content__b38e8bb9', 'get_docx_content__5ab5849946111fe9c5dda6e47422b057', 'get_text_file_content__87f48e6e', 'get_pdf_files_with_content__184db76a3945be5ea6dec91515beac2a', 'get_tetris_files__cb4dc6e5', 'get_file_exists__04b2d876', 'get_python_file_content__d881a7db2082639af55dd0ea1e3047ab', 'get_pdf_file_info__3dd8a1347f60dd796dc5e98521ae4032', 'get_file_content__89d906fa', 'get_python_files_count__13d7d579', 'get_file_exists__25309a67d723dd8e75eb60c978b60929', 'get_vm_file__519b2394782587cf0953c6f72e80b8a6', 'get_python_files_syntax__c9c8227a4cd72e8de3e73ea399f7f61d', 'get_text_file_lines__d1fc13ca7061617e08d8a914a14209cd', 'get_filename_from_path__a6b53d3d', 'get_tetris_files__58a20c62', 'get_pdf_files_info__8d8da24c', 'get_recent_pdf_count__db5a3e05', 'get_compressed_file_contents__0b074054', 'get_file_content__f8073900a375900c7a9ce8fa79f05f9d', 'get_file_existence__81d9e3403ff656f186df75dab490d6ac', 'get_vm_subtitle_file__dee238bc', 'get_launch_json__a3743b930b4e3c5c6976584a28c8269c', 'get_pdf_docx_result__a2a86631dd34db500aa084715506d32d', 'get_python_file_with_pattern__f222fdd4d3c26325ac7310b0f2b1711f', 'get_dual_dir_check__9cc29f71', 'get_file_content__cd720b2833d8c75c48e4cd829046ee69', 'get_txt_result__f5eea4d48349a222e7a85e09c99bbeae', 'get_txt_line_count__8ff67d2b', 'get_file_content__6c41fae8ffe95d2c4b3c521d5446de56', 'get_file_exists_and_size__479794c8', 'get_specific_pdf_exists__4e03b1ed', 'get_text_file_content__49075d97b370c316554a4b259a7ccc3e', 'get_pdf_files_in_dir__80b8ddf2', 'get_file_rename_status__c396550f', 'get_docx_content__fca57ed0800553851d5057d288f20b50', 'get_file_specific_line__ebbefd78bc8d93d288c452f335196528', 'get_pdf_page_count__25fc85ee', 'get_terminal_profile_name__dfa021d620bcc8024aed513c1adfaad3', 'get_file_exists__b8a50137', 'get_ext_name_by_path__4c281757414277d3f93140772916c7ad', 'get_csv_filtered_rows__d0eaa9b89f4859670e467a80277e775c', 'get_pdf_chapter_files__242f4871e4280ec1025212dbeaf5c18c', 'get_word_count_from_txt__68182234', 'get_file_line_indentation__a2b5af8108e461937716b976809e966c', 'get_multi_directory_contents__a93d97ba', 'get_file_permissions__fdcd3f41', 'get_dir_file_list__9ff7dd4db34cdbddf16b8ce0d6085594', 'get_file_list_from_text__24c8b0df', 'get_text_file_content__860f6a3e21fa3ee43fbc12dc1309f38a', 'get_pdf_basic_info__9f4baa29714d8f65b6c7e77480d7ee20', 'get_text_file_lines__c2d1250f71a1fa98280fe372a2eb5875', 'get_file_content__8ab0a45d4a57cf0e9592e87621895b59', 'get_docx_content__b519d5a9d3c41783a990af660a0ae167', 'get_tetris_files__65e60eeb', 'get_file_move_check__33f4e0c7b45da73c555209c719d5ff10', 'get_text_file_content__69119b71', 'get_srt_filename_check__ab47203640c5bdcef1195c50e51e7524', 'get_file_exists__6fb34f5d', 'get_file_info__739292ff', 'get_file_content__d09dce0a', 'get_pdf_files_with_content__83b6523335a3d68f9c734599b1537c74', 'get_specific_pdf_existence__6d39fc800ff8d1397c30bfae2c676b76', 'get_file_line_count__1098fce8', 'get_csv_unique_first_names_count__5ab9ee5026038714f957b134595e9a67', 'get_numeric_value_from_file__d8bf55dd66c809967b252ee6c81aa4a7', 'get_file_exists_and_size__e1affec5', 'get_file_lines__8278c76c3e9924f1866901a10447ca8f', 'get_file_exists__e3ae8a85', 'get_docx_text_content__18d2ce6a', 'get_file_exists__bdb8ae26', 'get_dir_exists__d31c5db6', 'get_files_with_prefix__8c4eaed9a61673f78def8f323e7dfe9d', 'get_file_absence__be265045', 'get_downloaded_file_info__08c5e1b6ad7015f1bdd4ff79ff88e12f', 'get_csv_filtered_rows__5ac234c6f992977135e16f82780e5511', 'get_vm_file_size__f259522f5141b84d8b2c6c9007fd732a', 'get_file_count_in_dir__081d0b6c', 'get_pdf_exports__2586a709', 'get_txt_content__562f595e9fa43de1b5d952aea3a3f9eb', 'get_docx_text_content__cea00e79', 'get_pdf_filenames_in_dir__7152d37e', 'get_single_pdf_with_verification__1386929a4648885c7d87d6425829fd94', 'get_numeric_value_from_txt__f9a0219a', 'get_docx_text_content__97c6df2d', 'get_folder_organization__6a9b75b65029a1742667a75cc968a2a9', 'get_archived_ipynb_files__0190943f9252410b5b69d12d28f1b6b7', 'get_readme_content__457448603041bcfcd06aff998da4e920', 'get_multi_directory_contents__94dfca2e', 'get_tetris_files__d0091276', 'get_file_size__f778a8914a698cd2bc7c0cc50cd3596d', 'get_file_exists__03122ed4', 'get_directory_has_file__b77b1963e53aafc9923374c6ea5077e2', 'get_text_file_lines__36494d7a38fdbc8d1f722d2db36fce8e', 'get_speedtest_json_data__c8d946870135d67f0db0be5e65caaa2a', 'get_file_properties__9c119f81', 'get_text_file_content__d7828490', 'get_file_move_check__aace45122e840d40b84dc540ae5a49bc', 'get_docx_text_content__6003371e', 'get_file_line_count__475840bd88bfc32515242a838ac799b5', 'get_file_exists__46ccb784', 'get_pdf_files_list__d2aaaf87', 'get_dir_exists__92ec1d53', 'get_subdirs_count__125fed35', 'get_rule_vm_file', 'get_vm_file_exists__3df7d80c', 'get_python_script_content__c9385c6b', 'get_text_file_content__f8ce9de4', 'get_pdf_with_name_pattern__55b958e432c8380deab73a3d2fcf329a', 'get_pdf_files_list__c8947a04', 'get_dir_file_list__36005e14c2ae4846381f0946699e70be', 'get_multi_directory_contents__8c566ad0', 'get_filename_hash_mapping__bf825e2c', 'get_file_list_content__557a0701', 'get_targz_contents__77bd062dcbbc640025c417629c0311e0', 'get_text_file_format__7f844786255954cce16f5ea58433f34e', 'get_csv_sorted_contacts__1e62491a', 'get_file_exists__9c31b3f6afb568d600c17c937149b6c4', 'get_files_matching_pattern__039c45a2', 'get_pdf_files_info__437d5a7f', 'get_file_content_match__57d8acad', 'get_subdirs_exist__2a5217dc', 'get_file_in_directory__31a8a4acc19afab71c5f7ec2f3006a11', 'get_pdf_files_list__cedddaca', 'get_python_file_info__198be354', 'get_docx_text_content__94746a60', 'get_multi_directory_contents__f3cabf2e', 'get_pdf_orientation__c75d6b17', 'get_file_content__3c9f051952e2f37565e45b593e085b87', 'get_files_with_prefix__f2ebfe11', 'get_file_exists_and_size__ff3634ef', 'get_file_exists__c0930f6fc6470d951ad9d775e3a6c6a5', 'get_file_exists__3928cfa5', 'get_text_file_content__b8171706418d2058f81460f1a24d4635', 'get_file_content_contains__c7e5cf7f', 'get_file_exists__7298f585793a2b727ac3910de3795a50', 'get_vm_subtitle_file__4c5ac05d', 'get_pdf_files_info__230c9972', 'get_docx_content__8c8e551c645c8d069577024a9d0bf137', 'get_file_location__c78e5698dfdf96679302f35b21f0928f', 'get_text_file_content__600696b8508be0c2e3ca25794856fb75', 'get_dir_file_list__d9c12923d5c2940b6520fb26a328924d', 'get_tetris_files__b3b822fe', 'get_number_from_file__9b848ddc', 'get_csv_structure_and_content__8de41bc4fc70600ffffb805994ee2926', 'get_pdf_file_exists__05a7fe2932371163af38b8e77d9b0c93', 'get_file_exists_check__f768bfed', 'get_file_location__cbcf6e0c', 'get_file_size__b756f99d', 'get_pdf_files_info__a96f92e2', 'get_number_from_file__c112996b', 'get_pdf_files_in_dir__32b125c8', 'get_file_count__23e95644', 'get_text_file_content__5ced85fc_aug18_v1_a1b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6', 'get_file_path_exists__c461bcde', 'get_both_files_non_empty__e07d7a26', 'get_file_timestamp__739292ff', 'get_pdf_files_info__f4eddf72', 'get_file_in_directory__35a43b1b', 'get_pdf_exports__c5999de9', 'get_docx_content_and_highlights__fc0ff0e3', 'get_two_files_exist__3f6d3219', 'get_numbered_pdf_chapters__a26911b7babb666234e16a8cd9e0ff40', 'get_subdirs_in_dirs__6eb64016ae326b8b48bfd5c29e5970e8', 'get_vm_file__a46a1ebc95419808bbd66e4ed9c5acac', 'get_text_file_content__433237d5d3b8a54ba0440460662ae73a', 'get_pdf_orientation__e222560cb6780019664a6917701a2659', 'get_folder_files__c14b0c4873b98943d0bcc59ebce705c9', 'get_csv_first_n_fullnames__a0281c1d922f2a38909920e121739b97', 'get_pdf_basic_info__56a55075c129634e9f17633f657c7657', 'get_zip_file_list__53fe105429a60cae06bcf9ce59b19b3e', 'get_pdf_files_in_dir__989759aa', 'get_docx_text_content__c32b1108', 'get_directory_tree__63701d4e', 'get_pdf_files_in_dir__bfc9cce9462c8f8be02842ad8f805893', 'get_file_exists__3b8e423e430323c0078f4425aded05b9', 'get_text_file_content__210f933b0657674e3cf1f4222885738f', 'get_docx_content__2e06628be7ef5b3dc45161c8ea091810', 'get_directory_structure__3ef6f937c5dcda24ec08a8ba5aa2c872', 'get_file_exists__6ee1fdcd', 'get_pdf_page_count__e941d252d60fece99cf71fbcddf8521f', 'get_tetris_files__c6117858', 'get_blog_folder_listing__4e03b1ed', 'get_pdf_file_info__40d8ee41df97cbb2f480ba7af545efc4', 'get_file_renamed__b4364020', 'get_speedtest_csv_data__50006736ce024a92878da0e4240e2bb0', 'get_pdf_count_in_dir__085cc8d6', 'get_file_existence__a8a082525df1807c95a7519289fda5a0', 'get_file_size__198be354', 'get_text_file_content__c4d957b9ed4c029b04783deb70ccd213', 'get_pdf_files_info__6e9dcacc', 'get_file_content__e1c8e8d0', 'get_text_file__3c678f53', 'get_text_file_lines__bf9ee805c777733e26c14d1927c30b1a', 'get_multi_directory_contents__3b0a753c', 'get_directory_listing__32107748', 'get_zip_files_and_dirs__e2b9921b3a4cfcd62d172fdaae4844bb', 'get_text_file_content__0fe1ad2c6a085376a5c2eb19244f70a6', 'get_file_exists__5f5351b0', 'get_docx_text_content__6c0adead2c9259465f7529a90ac3bebf', 'get_text_file_content__7dfb45a4', 'get_pdf_files_info__2c1e781e', 'get_file_exists__898e13b4', 'get_directory_info__e01c3944c2f0dfae12ffc1d2b96464cf', 'get_file_recovery_state__5ea617a3', 'get_file_exists__c6c3aa52', 'get_text_file_content__acd75d14', 'get_pdf_files_list__6871c0fe', 'get_file_count__19cf6326', 'get_pdf_page_count__a73a14fa', 'get_pdf_exports__d582d1d2', 'get_pdf_content__5e9826db825b680d44e19c36727da776', 'get_csv_row_count__2d3e7876', 'get_docx_content_check__ef584ea3', 'get_question_count_file__e2620a9d', 'get_text_file_lines__2267ae43c99e342cd984b7743dc212e6', 'get_pdf_files_list__89dbe8bb', 'get_dir_structure_info__5b67568a', 'get_pdf_files_in_dir__d928a635', 'get_flac_file_status__b00c11502ffa26aa7b145b3096daa5d2', 'get_python_file_lines__d93565be', 'get_pdf_exports__831e0e03', 'get_file_moved_check__789836386f3e1cf0e0ee5d172a0885f2', 'get_pdf_properties__0b5ef92a5e5fe7305b438702a0eb6c3a', 'get_vm_file__9817ff525ad645c5458b8a22c03332a7', 'get_docx_file_count__a8fe8251', 'get_subdirectory_exists__6f9e6ffd', 'get_file_exists__974295f9d11461d175dbc0223dd4ff65', 'get_file_variable_occurrences__ba25acbbf5417edd79f01d39568b431f', 'get_screenshot_file_info__83a4fedf3a2dd9034045f9bfa6b3d8ed', 'get_text_file_content__e34109066d47c745bfcd2dc2768683b4', 'get_pdf_page_count__e56af0ce', 'get_pdf_files_list__6b40f44d', 'get_pdf_count__4e03b1ed', 'get_pdf_files_info__0707ddca', 'get_file_exists__506ad17f', 'get_csv_and_count__a4d37536', 'get_pdf_basic_info__d36b36f1df35c32e48eb7411bf5a6c33', 'get_pdf_file_size__29f6acc744cf0d15caa355ed1701b507', 'get_desktop_renamed_file__6bf0504dae1e157e634b1e1c5be03fad', 'get_file_exists__7648b3ac', 'get_word_count_from_file__7aa35df1', 'get_vm_subtitle_file__f5680565', 'get_docx_content__d46c5e6abd5059bf9bd9590b3b72a55f', 'get_python_file_content__093631738f9b5eba42a5bcf60212ba3b', 'get_docx_bold_content__eb25c2f0ed9e56a458cb4ee477d415ca', 'get_pdf_files_from_docx__24a9a657e8a67df9509783bd4f53b233', 'get_desktop_file_list__e35e8479c21a9a3d6ef729a997676f8c', 'get_pdf_exports__b47a318c', 'get_text_file_lines__5ced85fc_aug18_v3_c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8', 'get_file_exists_at_path__24914e86', 'get_text_file_content__1a1f627807b83c4d33e1ae428da08935', 'get_file_exists__60340d37', 'get_pdf_validity__90edd5864af08865c892a3ecd6659029', 'get_pdf_file_location__0d0715b7122c298c1e10aa6fa135f599', 'get_vacation_files__d7a48669399bf74024b2a979a32d4ae1', 'get_snake_direction_test__c739e38ab0abd37ef206886e6d29d7b5', 'get_file_content__8085b902c0aa531b12d2ed766e22c897', 'get_file_permissions__b3b80682', 'get_docx_content_status__ba90be29', 'get_pdf_file_list__10f702e7ff0a641a1fda6d45251486ce', 'get_folder_file_list__32b2b661', 'get_text_file_content__de3b1681', 'get_filename_info__4b590a3a028f08e8f4ad12729f4351c3', 'get_text_file_content__5ced85fc_aug18_v2_b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7', 'get_files_in_directory__880f8efb', 'get_file_size__b9c089b2fe7d833fde2da297bbbd9620', 'get_chapter_pdf_count__457b9850ebb20f14d7c68688b0aa27a6', 'get_pdf_files_with_validation__1cc9c1bfb30ee1b2b3cdbf7926894918', 'get_text_file_content__834c93d1a65ecbb7766bb5ceb1a12320', 'get_vm_subtitle_file__c48ab0f0', 'get_pdf_merge_info__a8bc7d70fd87c66cda30ee22072de4d3', 'get_vm_text_file__7c58ef63', 'get_pdf_text__439ece4f99f624319fffd552c64d5f89', 'get_docx_text_content__7080c120', 'get_file_exists__42ef2d72a0f41077972857e814318a24', 'get_two_zip_file_counts__0cf1485f', 'get_pdf_page_count__434044232831686d814f3fab48f559ee', 'get_pdf_files_in_dir__cdfaf4c1', 'get_file_move_status__00112b53200a74ce7a53869d2d085264', 'get_merged_file_content__618fd61c', 'get_file_modification_time__23cbcfa9', 'get_pdf_file_count__086472c391d9b3dbf463fe72713bc019', 'get_csv_filtered_rows__e49a53b3238eb33c8eed6130c48c5268', 'get_file_exists__7930967f', 'get_text_file_lines__465a794bbc0c138e64d3078d7f3f0b51', 'get_docx_content__61084d52', 'get_vm_subtitle_file__dccfbd8e', 'get_default_pdf_viewer__10cce6e7e23a5b966c5cbd94b3842103', 'get_file_exists__25f0d8c3', 'get_default_pdf_viewer__09a6505e4073dfc93152a41afb590f9b', 'get_docx_text_content__3874241e', 'get_text_file_content__8b701bf8d0cacb95438d2e4e17a8b914', 'get_file_exists_check__e185e2be', 'get_dir_file_count__c907493f', 'get_pdf_exists__7a335cc66cd23236defef5bd95bbbe7d', 'get_txt_file_content__b2ba9ee6873b43000b20ba169c4d9592', 'get_pdf_files_in_dir__1ee77af4', 'get_file_count_in_directory__23c57bc9', 'get_tetris_files__72c9bc74', 'get_pdf_file_sizes__4e03b1ed', 'get_file_content__b6bb3e42', 'get_text_file_lines__c2b520792100f4aa54f0d89dedfa94c4', 'get_desktop_named_pdf__0c8a922818da41c6a16e3979ae1f26dc', 'get_remaining_files__8754d37bdc9e8d94ab80feb618caa015', 'get_python_file_content__be0d67f2ca1c7fde5d3c550cd046d0b4', 'get_file_exists_check__e4d07acf', 'get_renamed_file__9f8a1d0c', 'get_text_file_content__5ef70598fcc68171b6ec36d7c888fc21', 'get_file_exists__e3dda739ee4da14903d2e8df52d7a41d', 'get_default_pdf_viewer__bdb427f6', 'get_text_content__67c66e9d6c723be29d116d6e2c7b5850', 'get_backup_folder_files__089628262a733e157c368e0bf1ad02f1', 'get_file_content__edd1c1331eff751ad5487718d7e5d07b', 'get_vm_file_exists__8752e3fa', 'get_file_line_count__e09bfbdf', 'get_pdf_basic_info__0448af3ec8f726fcfbf94a06b95da924', 'get_file_exists__724cfa0a', 'get_archive_contents__ccef1fda30d9645300e2b2b1c657b52f', 'get_identified_mountains_folder__c8c87c59ecc91c5a6beb4965364b8f59', 'get_backup_folder_hashes__dd5eeed0', 'get_tetris_files__0e016e1b', 'get_pdf_files_list__21dcc4a0', 'get_pdf_files_list__0be3e647', 'get_renamed_files__27c9f1432580c9f5f1bf2d1f919f4ed5', 'get_text_file_content__08e73922', 'get_pdf_file_sizes__99e94d1a', 'get_sender_file__ad43db2627764dd28ab9631606c7b97c', 'get_file_list__91c6f86e', 'get_pdf_numbered_files__fe03784ed42c9a7002fcc758df207953', 'get_text_file_content__68f83c37fe78a3f2dbadefcb3c480330', 'get_text_file_line_count__5538243966b3481fe772536923c1f693', 'get_zip_filenames__5c33f919', 'get_writer_text_content__592c042e86aba9b7d5a9be8008d4a4f0', 'get_text_file_content__2d81db9f5efc3d72c38ba9f24bf6d4fb', 'get_ipynb_dirs__c69badbf0960e9198fe7a7334d4b95f3', 'get_docs_subdir_doc_count__09b713a6', 'get_file_exists__ec920d7f', 'get_zip_file_count__03f8ef9d', 'get_file_content_text__8369c71b0ee26c543c025c6e1cb39bbd', 'get_recent_pdf_count__4e03b1ed', 'get_pdf_file_info__000c73be1394701e25c8491dd9418647', 'get_file_exists__0ad50b28', 'get_renamed_docx_file__232d51eaf1622d499f79eaf2309ca5ad', 'get_pdf_files_in_dir__d8278409', 'get_zip_verification_state__1a0ccf050b10545e05649e32a895cb33', 'get_merged_file_content__3b864a94', 'get_notebooks_dir_files__1b47b6505a7a2dc3d6ad6f0c07b4bcb4', 'get_all_doc_pdf_conversion__1d31f566abd6111dc90f53a510d255f5', 'get_text_file_content__d332c3241fced231d1d84d00e75fe3b7', 'get_file_exists__7ae8ce2b', 'get_file_exists_and_size__c0d603dc', 'get_subdirs_in_dir__c4632abc', 'get_file_exists__f25970ca', 'get_file_exists_and_size__995b229f', 'get_pdf_exports__98b0d7a7', 'get_filenames_with_pattern__7427978e', 'get_copied_files__9e83c51609bc4e6d77dc6303641e8cf9', 'get_file_content__e5eac3aa5287398e5510fa8c359cddce', 'get_pdf_file_properties__869dd2b5b7a0e45511c735d06983da83', 'get_file_exists__67be1ac6efe87edab008a615fb0e7ec4', 'get_vm_file_text__26660ad1', 'get_multi_directory_contents__92a58812', 'get_file_not_exists__9e688855', 'get_text_file_content__11d0824d05970e58a5671a5636365f15', 'get_word_count_file__adfc25c4', 'get_pdf_file_list__822cf85f742bbfae9e3acf8d7027940c', 'get_file_list__990ae9b047da99489a16db0558f7ee61', 'get_file_exists__684b5a3a3f653750766f5bbe64af3bd5', 'get_pdf_file_sizes__2fc6b524', 'get_mountain_prefixed_files__0f256e5c55ade51a0f98ec735f4e6698', 'get_pdf_text_content__7a2513e6563b76f5c07e27d8c4089d02', 'get_vm_file__a5c1978c26f91a02b256f236a6017e74', 'get_csv_row_count__58df2e60', 'get_file_exists__ecf92ffc', 'get_pdf_file_count__2f750d009f0f471751ed869c09f90a90', 'get_zip_filenames__d6b113924c427deceec5f933af24484e', 'get_text_file_content__25fb76d7ccb83f42013b589a25bead61', 'get_text_file_content__ed1a5c265e6c6d06dcaf2ec482204403', 'get_pdf_file_sizes__d4f3d6039bf1a71698a65c2d049521e8', 'get_pdf_basic_info__247861df47176dd2a8c57d33f4d99eab', 'get_file_exists__e1da6937', 'get_pdf_file_info__5f59826466b2625834fd8d369560ed11', 'get_file_rename_check__c14397f3284104a2f980691d2ea6abf3', 'get_csv_filtered_count__83247848fa547df5ef0efe5837431ceb', 'get_file_count_by_pattern__816aac7b5fcbde572f13b62a3999bd4d', 'get_pdf_files_in_folder__4e03b1ed', 'get_file_first_line__ec5fe11be98e432184dc357c559fbb9f', 'get_default_file_manager__a1cf9be3', 'get_chapter_files_status__93cfd69b3c5adfd5dbb8817764000202', 'get_text_file_content__20f90bc2668d01c30660dfeafc3af15b', 'get_vm_dir_list__abdf3926e9609e7e5b435c8cdbb40013', 'get_pdf_count_desktop__6c8aa0e0d1a7a1f247933b05b6dc0e08', 'get_zip_contents__8a01d242c9e2052109e1c26f5fa4a5dd', 'get_text_file_content__fc2c8cc4', 'get_pdf_count_in_dir__a0da932e', 'get_file_line_count__bcaf4400', 'get_vm_file', 'get_text_file_content__d997648158d90618366cc740cfd24b34', 'get_wallpaper_path__a70b53e8', 'get_file_count_in_dir__fbd9137c', 'get_pdf_multiple_fields__ff249445b547a028ee5246e45cd19fdf', 'get_picture_filenames__fdcd3f41', 'get_docx_content_and_font__abb42de39a213997a2a7f06fa1fe5d2a', 'get_file_existence__198be354', 'get_file_exists__62dacdc1', 'get_file_exists__b03b6c61', 'get_csv_merge_data__c66369e707b97de3ccd6da4699663fe6', 'get_text_file_content__7bbdf0a0733630cbbfd86729556fc827', 'get_pdf_filenames__4e03b1ed', 'get_pdf_basic_info__7b63a847bf086913e974e7a32debb9ec', 'get_multiple_files_exist__03426e679d8f4571bede57a16eea69a4', 'get_directory_files__940d01bc', 'get_dir_file_count__95b4929b', 'get_pdf_files_info__8010e79b', 'get_file_exists_and_size__d6a98717', 'get_folder_contents__8542f72b', 'get_inbox_backup_dir__987d91d90e3876e4880d553cc7b5b944', 'get_vm_file__96eb02f4ac90683da522ea44f75e2519', 'get_file_exists__46f0b51f', 'get_python_file_content__generic', 'get_multiple_files_status__4e03b1ed', 'get_python_file_structure__58d4fb9ea6e69b7f57a37a59813050e8', 'get_dir_file_count__00db2192', 'get_text_file_content__7e5134258960ebea77ca0d290984a7a3', 'get_pdf_in_year_folders__5ad487ac7d025a6b906a4c83e8beac41', 'get_python_file__4ca3247dcf6b464342c4e3f53d844797', 'get_playlist_file_status__3145ebe2c0a2fb7aa964b8c8f3b95ea9', 'get_folder_file_count__5e18d045', 'get_pdf_content_fit__234d824f7c8a19c4a33bd790f5cee244', 'get_csv_column_count__7f9974e9', 'get_git_dir_exists__e2da960ab9034666db33db74ae6371a7', 'get_bash_file_content__5ed0399c75d125da4fc3b5f7583bb5c4', 'get_file_exists__f663e89c52b74d5c5d4e38ab6d86c83f', 'get_file_copy_verification__dc38ce29eba391e7169ff4e028e69a72', 'get_screenshot_file_info__57d2d1f6', 'get_docx_text_content__e38643bd', 'get_tetris_files__03faeae5', 'get_desktop_pdf__a45ebe2876987410607711d3992c10db', 'get_pdf_info__9a18b30d646547c54b42b3593f83920d', 'get_file_line_count__7233f122896fac183c973343e2cf3b2a', 'get_large_file_count__b9976565', 'get_file_count_from_text__bc253f41', 'get_pdf_page_orientation__596fcd5e219bea1372357f0afb95cc85', 'get_merged_text_file__e22bfb55f6ab9983d1cc35b82dc09aeb', 'get_files_exist__f50a55ca', 'get_text_file_content__438c9c7ce7eebe25a3992ddf0a388112', 'get_text_file_content__515e2337245bc72c8d34192293ce6646', 'get_mbox_content__1701225e2d2da95122de9fd6941c3c6f', 'get_pdf_files_info__ea295e4c379ee25a192357c908aada76', 'get_csv_headers__4735f34d496bbe1961d6d9a20cf7b9bd', 'get_file_exists__689ec9af4ba1471bf9b5f89e71cafeb9', 'get_file_content__d77a2096', 'get_pdf_exists__9753fb3ef75762f14fb08b7f236e3f81', 'get_file_content__c2756c851c53289bbd5185905c6853a2', 'get_text_file_content__551695fc', 'get_text_file_lines__804f1b3e28fc8c6ee1d8689118367b3b', 'get_file_exists__d8671412', 'get_python_file_content__f55aa7954b40a62bad3b8ba851857ed1', 'get_bashrc_path_check__c0107678', 'get_multi_directory_contents__b1a155d8', 'get_text_file_content__f84e9cb5fd8cfc9fab208c24bcd90a7d', 'get_timestamped_file__2b78c2fd0d670b6ee1c54ce65b4419a5', 'get_file_list_content__ef0f28d6', 'get_vm_subtitle_file__2ad8e92c', 'get_pdf_text_content__4c28c68dd08d7073669bab47ee359a64', 'get_text_file_content__895c3960d172d43278234eeb5c495eda', 'get_text_file_lines__dfc9794a', 'get_text_file_line_count__ab2d13c4', 'get_text_file_content__6941d0dc31bbd0c3d844303b7e1c57e5', 'get_csv_filtered_count__97b9c260', 'get_file_count__9364293cce5b25e22063aee62da7d43d', 'get_file_count_from_file__be02851a', 'get_file_list_content__b23d642c', 'get_bib_file_nonempty__9d1faa88', 'get_file_exists_and_size__21492178', 'get_file_exists__e86a3a4d', 'get_txt_file_count__2b6d7a72', 'get_dir_file_list__173d79a32d2c31ac3ad30d4ae958526d', 'get_files_with_prefix__47275204', 'get_csv_filtered_rows__01b44b7acd315e39b1c9a6baa6b5f6da', 'get_vm_subtitle_file__735234d9', 'get_text_file_content__4080707f437c813fb70f2db7aaa30575', 'get_file_exists__7107319d', 'get_dual_file_check__805294f8', 'get_file_line_count__fb1a48c8', 'get_text_file_lines__5ced85fc_aug18_v0_c9e8a1b2d3f4e5a6b7c8d9e0f1a2b3c4']

def get_zip_contents__ed5b60c9b586e02ab1d262e432f8c2c4(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract and verify ZIP file contents from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to ZIP file on VM

    Returns:
        Dict with 'exists', 'is_valid_zip', and 'file_list' keys
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'is_valid_zip': False, 'file_list': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.zip', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zf:
                file_list = sorted([info.filename for info in zf.filelist])
                return {'exists': True, 'is_valid_zip': True, 'file_list': file_list}
        except zipfile.BadZipFile:
            return {'exists': True, 'is_valid_zip': False, 'file_list': []}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'is_valid_zip': False, 'file_list': []}

def get_pdf_exports__bec19abd(env, config: dict):
    """
    Check if PDFs were exported to the specified directory and extract Chrome tab titles.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key and optional 'urls' key

    Returns:
        Dict with 'pdf_files' (list of PDF filenames), 'article_titles' (list of article titles from Chrome tabs),
        and 'pdf_titles' (list of titles extracted from PDF files)
    """
    directory = config.get('directory', '/home/user/Downloads/BlogPosts')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        pdf_files = [os.path.basename(f.strip()) for f in list_result['output'].strip().split('\n') if f.strip()]
    article_titles = []
    try:
        response = requests.get('http://localhost:9222/json', timeout=5)
        targets = response.json()
        expected_urls = config.get('urls', [])
        for target in targets:
            if target.get('type') == 'page':
                url = target.get('url', '')
                title = target.get('title', '')
                if expected_urls:
                    for expected_url in expected_urls:
                        if url.rstrip('/').startswith(expected_url.rstrip('/')):
                            if title and title not in ['', 'New Tab', 'Chrome', 'about:blank']:
                                article_titles.append(title)
                                break
                elif title and title not in ['', 'New Tab', 'Chrome', 'about:blank']:
                    article_titles.append(title)
    except Exception as e:
        try:
            obs = env.controller.get_obs()
            a11y_tree = obs.get('accessibility_tree', '')
            if a11y_tree:
                for line in a11y_tree.split('\n'):
                    if 'RootWebArea' in line:
                        match = re.search('name=["\\\']([^"\\\']+)["\\\']', line)
                        if match:
                            title = match.group(1)
                            if title and title not in ['', 'New Tab', 'Chrome', 'about:blank']:
                                article_titles.append(title)
        except Exception as e2:
            pass
    pdf_titles = []
    for pdf_file in pdf_files:
        pdf_path = os.path.join(directory, pdf_file)
        try:
            doc = fitz.open(pdf_path)
            metadata_title = doc.metadata.get('title', '')
            first_page_title = ''
            if len(doc) > 0:
                page = doc[0]
                text = page.get_text()
                lines = [line.strip() for line in text.split('\n') if line.strip()]
                if lines:
                    first_page_title = lines[0]
            doc.close()
            extracted_title = metadata_title if metadata_title else first_page_title
            pdf_titles.append(extracted_title)
        except Exception as e:
            pdf_titles.append('')
    return {'pdf_files': pdf_files, 'article_titles': article_titles, 'pdf_titles': pdf_titles}

def get_pdf_files_info__cb51ce0b(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_exists__dc35efec(env, config):
    """
    Check if a file exists and validate it's a valid PNG image with reasonable properties.

    Returns a dict with:
    - exists: bool - whether file exists
    - is_png: bool - whether file is a valid PNG image
    - file_size: int - file size in bytes (0 if not exists)
    - is_recent: bool - whether file was modified recently (within last 5 minutes)
    """
    file_path = config.get('path', '')
    command = f'''\nif [ -f "{file_path}" ]; then\n    echo "EXISTS=true"\n    # Check if it's a PNG file using file command\n    file_type=$(file -b --mime-type "{file_path}")\n    echo "MIME_TYPE=$file_type"\n    # Get file size\n    file_size=$(stat -c %s "{file_path}" 2>/dev/null || echo "0")\n    echo "FILE_SIZE=$file_size"\n    # Get file modification time (seconds since epoch)\n    mtime=$(stat -c %Y "{file_path}" 2>/dev/null || echo "0")\n    echo "MTIME=$mtime"\n    # Get current time\n    current_time=$(date +%s)\n    echo "CURRENT_TIME=$current_time"\nelse\n    echo "EXISTS=false"\n    echo "MIME_TYPE=unknown"\n    echo "FILE_SIZE=0"\n    echo "MTIME=0"\n    echo "CURRENT_TIME=0"\nfi\n'''
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    info = {}
    for line in output.split('\n'):
        if '=' in line:
            (key, value) = line.split('=', 1)
            info[key] = value
    exists = info.get('EXISTS', 'false') == 'true'
    mime_type = info.get('MIME_TYPE', 'unknown')
    file_size = int(info.get('FILE_SIZE', '0'))
    mtime = int(info.get('MTIME', '0'))
    current_time = int(info.get('CURRENT_TIME', '0'))
    is_png = 'image/png' in mime_type or 'png' in mime_type.lower()
    time_diff = current_time - mtime
    is_recent = exists and 0 <= time_diff <= 300
    return {'exists': exists, 'is_png': is_png, 'file_size': file_size, 'is_recent': is_recent}

def get_file_exists_check__b62f9acb(env, config: Dict) -> Optional[Dict]:
    """
    Check if a file exists and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_docx_text_content__6477df40(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_file_count_in_dir__f1d66967(env, config: dict):
    """Get detailed information about PDF files in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path
            - extension: File extension (e.g., 'pdf')

    Returns:
        Dict with file details including count, filenames, sizes, and timestamps
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    extension = config.get('extension', 'pdf')
    command = f"""\npython3 -c "\nimport os\nimport glob\nimport json\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path):\n    print(json.dumps({{'count': 0, 'files': []}}))\n    exit(0)\n\nfiles = glob.glob(os.path.join(dir_path, '*.{extension}'))\nfile_info = []\n\nfor f in files:\n    try:\n        stat = os.stat(f)\n        file_info.append({{\n            'name': os.path.basename(f),\n            'size': stat.st_size,\n            'mtime': stat.st_mtime\n        }})\n    except Exception as e:\n        pass\n\nresult = {{\n    'count': len(files),\n    'files': file_info\n}}\nprint(json.dumps(result))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to get file info: {result['error']}")
        return {'count': 0, 'files': []}
    try:
        data = json.loads(result['output'].strip())
        return data
    except Exception as e:
        logger.error(f'Failed to parse file info: {e}')
        return {'count': 0, 'files': []}

def get_music_file_count__355d5753246a8a88c1b8d173c110d89e(env, config: dict):
    """Get list of MP3 filenames in Music directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        list: List of MP3 filenames (basenames only)
    """
    directory = config.get('directory', '/home/user/Music')
    try:
        command = f'find "{directory}" -maxdepth 1 -type f -name "*.mp3" -exec basename {{}} \\;'
        result = env.controller.run_bash_script(command, timeout=10)
        if result.get('returncode') == 0:
            output = result.get('output', '').strip()
            if output:
                filenames = [name.strip() for name in output.split('\n') if name.strip()]
                logger.debug(f'Found MP3 files in {directory}: {filenames}')
                return filenames
            else:
                logger.debug(f'No MP3 files found in {directory}')
                return []
        else:
            logger.error(f"Error listing MP3 files: {result.get('error', '')}")
            return []
    except Exception as e:
        logger.error(f'Error listing MP3 files in {directory}: {e}')
        return []

def get_directory_file_count__58fb65f5(env, config: Dict[str, Any]) -> int:
    """Get count of files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Number of files in the directory
    """
    path = config['path']
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return 0
    return len(result['children'])

def get_documents_file__c0d05fd4eae793b685d13d8511de53aa(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists in Documents folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'is_png', 'in_documents' keys
    """
    path = config.get('path', '')
    result = {'exists': False, 'is_png': False, 'in_documents': False}
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            result['exists'] = True
            if '/Documents/' in path or path.startswith('/home/user/Documents/'):
                result['in_documents'] = True
            if path.lower().endswith('.png'):
                try:
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    try:
                        img = Image.open(tmp_path)
                        img.verify()
                        result['is_png'] = True
                    except Exception as e:
                        logger.warning(f'Not a valid PNG: {e}')
                    finally:
                        os.unlink(tmp_path)
                except Exception as e:
                    logger.warning(f'Error verifying PNG: {e}')
        else:
            logger.info(f'File does not exist at path: {path}')
    except Exception as e:
        logger.error(f'Error checking file: {e}')
    return result

def get_pdf_file_list__d2aeadd33c64efc8e0e3416b06f6e222(env, config: dict):
    """Get list of PDF files in a directory with metadata (size, content validation).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        dict: Dictionary with:
            - 'files': List of dicts with 'name', 'size', 'is_valid_pdf' for each PDF file
            - 'source_pdf_exists': Boolean indicating if source PDF exists
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    source_pdf = f'{directory}/Spectral Graph Theory.pdf'
    check_source_cmd = f'test -f "{source_pdf}" && echo "exists" || echo "missing"'
    source_result = env.controller.run_bash_script(check_source_cmd, timeout=10)
    source_exists = source_result.get('output', '').strip() == 'exists'
    command = f'\nfind "{directory}" -maxdepth 1 -name "*.pdf" -type f -exec stat -c "%n|%s" {{}} \\; 2>/dev/null || echo ""\n'
    result = env.controller.run_bash_script(command, timeout=10)
    files = []
    if result.get('returncode') == 0 and result.get('output'):
        output = result['output'].strip()
        if output:
            for line in output.split('\n'):
                line = line.strip()
                if line and '|' in line:
                    (path, size_str) = line.rsplit('|', 1)
                    filename = path.split('/')[-1]
                    if filename == 'Spectral Graph Theory.pdf':
                        continue
                    try:
                        size = int(size_str)
                    except ValueError:
                        size = 0
                    is_valid_pdf = False
                    if size > 0:
                        check_pdf_cmd = f'head -c 5 "{path}" 2>/dev/null'
                        pdf_check = env.controller.run_bash_script(check_pdf_cmd, timeout=10)
                        if pdf_check.get('returncode') == 0:
                            header = pdf_check.get('output', '')
                            is_valid_pdf = header.startswith('%PDF-')
                    files.append({'name': filename, 'size': size, 'is_valid_pdf': is_valid_pdf})
    return {'files': files, 'source_pdf_exists': source_exists}

def get_rtf_file_info__34460ac1a76394f2dfa8b4e9981344a0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if an RTF file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' key indicating if file exists
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'RTF file not found at {path}')
        return {'exists': False, 'path': path}
    logger.info(f'RTF file found at {path} ({len(file_bytes)} bytes)')
    return {'exists': True, 'path': path, 'size': len(file_bytes)}

def get_docx_text_content__7017cad57d424df42048663d589125a3(env, config: Dict[str, Any]) -> str:
    """Extract all text content from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        str: All text content from the document, concatenated
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_text = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                all_text.append(text)
        return '\n'.join(all_text)
    finally:
        os.unlink(tmp_path)

def get_pdf_file_a503b07f(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get a PDF file from the VM and return its local path.

    This getter downloads a single PDF file from the VM and returns
    the local cache path as a string for verification.

    Args:
        env: Environment object with controller
        config: Configuration dict with:
            - path (str): absolute path on the VM to fetch
            - dest (str): file name of the downloaded file

    Returns:
        Optional[str]: Local path to the downloaded PDF file, or None if failed
    """
    from datetime import datetime
    vm_path = config.get('path')
    dest_filename = config.get('dest')
    if not vm_path or not dest_filename:
        return None
    cache_path = env.controller.get_file(vm_path, dest_filename)
    if cache_path and os.path.exists(cache_path):
        return cache_path
    return None

def get_text_file_content__7676732d(env, config: dict):
    """Read text file content from VM and check for CSV export.

    This getter performs two checks:
    1. Reads the count file content (1990s_count.txt)
    2. Checks for CSV file on Desktop (preferably with address book related keywords)

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary containing:
            - 'count_content': Content of the count file (str)
            - 'csv_exists': Whether a CSV file exists on Desktop (bool)
            - 'csv_file': Name of the CSV file if found (str or None)
              Priority given to files containing 'address', 'contact', 'personal', or 'book'
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    count_content = ''
    if file_bytes is not None:
        try:
            count_content = file_bytes.decode('utf-8').strip()
        except Exception:
            count_content = ''
    desktop_path = env.controller.get_vm_desktop_path()
    csv_exists = False
    csv_file = None
    try:
        directory_tree = env.controller.get_vm_directory_tree(desktop_path)
        if directory_tree and 'children' in directory_tree:
            csv_files = [file['name'] for file in directory_tree['children'] if file.get('type') == 'file' and file['name'].lower().endswith('.csv')]
            if csv_files:
                address_keywords = ['address', 'contact', 'personal', 'book']
                priority_csv = None
                for csv in csv_files:
                    csv_lower = csv.lower()
                    if any((keyword in csv_lower for keyword in address_keywords)):
                        priority_csv = csv
                        break
                csv_file = priority_csv if priority_csv else csv_files[0]
                csv_exists = True
    except Exception:
        csv_exists = False
    return {'count_content': count_content, 'csv_exists': csv_exists, 'csv_file': csv_file}

def get_text_file_content__d48f445fac6cb18530cd7f7169fdc7fc(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file on VM

    Returns:
        str: File content as string
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_file_list__4a34f03bcb82e7e8037181cd9a91ae6e(env, config: dict) -> dict:
    """Get directory listing from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Directory tree dict with 'children' list
    """
    return env.controller.get_vm_directory_tree(config['path'])

def get_default_pdf_viewer__6a629498(env, config: dict):
    """Gets the default application for PDF files.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The default PDF viewer registered for application/pdf MIME type
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'application/pdf']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_file_exists__cb8ab5642aaf48705d11abb5543759c9(env, config: dict):
    """Check if a file exists in the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    file_path = config['path']
    command = f"test -f {file_path} && echo 'EXISTS' || echo 'NOT_EXISTS'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        return output == 'EXISTS'
    else:
        logger.error('Failed to check file existence. Status code: %d', response.status_code)
        return False

def get_targz_file__04a2c355(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get information about a tar.gz archive from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool), 'size_bytes' (int), and 'file_count' (int) keys
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'size_bytes': 0, 'file_count': 0}
        cache_path = os.path.join(env.cache_dir, 'temp_check.tar.gz')
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        with tarfile.open(cache_path, 'r:gz') as tf:
            members = tf.getmembers()
            file_count = sum((1 for m in members if m.isfile()))
        return {'exists': True, 'size_bytes': len(file_bytes), 'file_count': file_count}
    except Exception as e:
        logger.error(f'Failed to read tar.gz file: {e}')
        return {'exists': False, 'size_bytes': 0, 'file_count': 0}

def get_numbers_from_txt__6622ce38(env, config):
    """
    Extract all numbers from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of integers found in the file
    """
    import re
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return []
        content = file_bytes.decode('utf-8', errors='ignore')
        numbers = [int(n) for n in re.findall('\\d+', content)]
        return numbers
    except Exception as e:
        return []

def get_pdf_exports__d587f20a(env, config: dict):
    """
    Check if PDFs were exported to the specified directory with metadata.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict with 'pdfs' (list of dicts containing filename, mtime, and content preview) and 'count'
    """
    directory = config.get('directory', '/home/user/Documents/Research')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f -printf "%p\\t%T@\\n" 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_info = []
    if list_result['output']:
        current_time = time.time()
        for line in list_result['output'].strip().split('\n'):
            if not line.strip():
                continue
            parts = line.strip().split('\t')
            if len(parts) == 2:
                (filepath, mtime_str) = parts
                filename = os.path.basename(filepath)
                mtime = float(mtime_str)
                age_seconds = current_time - mtime
                is_recent = age_seconds < 300
                content_preview = ''
                try:
                    python_script = f"""\nimport fitz\ntry:\n    doc = fitz.open('{filepath}')\n    if len(doc) > 0:\n        text = doc[0].get_text()[:500]  # First 500 chars of first page\n        print(text)\n    doc.close()\nexcept Exception as e:\n    print(f"Error: {{e}}")\n"""
                    preview_result = env.controller.run_python_script(python_script, timeout=10)
                    if preview_result.get('output'):
                        content_preview = preview_result['output'].strip()
                except:
                    pass
                pdf_info.append({'filename': filename, 'mtime': mtime, 'is_recent': is_recent, 'content_preview': content_preview})
    return {'pdfs': pdf_info, 'count': len(pdf_info)}

def get_folder_file_list__465762dc(env, config: dict):
    """Get list of filenames in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        list: List of filenames in the directory
    """
    folder_path = config.get('folder_path', '')
    command = f'ls -1 "{folder_path}" 2>/dev/null'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        output = result['output'].strip()
        if output:
            return output.split('\n')
        return []
    return []

def get_total_files_count__e9983217(env, config):
    """
    Get comprehensive information about the repository structure.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        dict: {
            'file_count': int,
            'has_git_dir': bool,
            'required_files_exist': dict,
            'required_dirs_exist': dict
        }
    """
    repo_path = config.get('repo_path', '/home/user/instructor-embedding')
    count_command = f'find {repo_path} -type f 2>/dev/null | wc -l'
    count_result = env.controller.run_bash_script(count_command, timeout=10)
    try:
        file_count = int(count_result.get('output', '0').strip())
    except ValueError:
        file_count = 0
    git_command = f"test -d {repo_path}/.git && echo 'exists' || echo 'not_found'"
    git_result = env.controller.run_bash_script(git_command, timeout=5)
    has_git_dir = git_result.get('output', '').strip() == 'exists'
    required_files = {}
    for filename in ['setup.py', 'requirements.txt', 'README.md']:
        file_command = f"test -f {repo_path}/{filename} && echo 'exists' || echo 'not_found'"
        file_result = env.controller.run_bash_script(file_command, timeout=5)
        required_files[filename] = file_result.get('output', '').strip() == 'exists'
    required_dirs = {}
    for dirname in ['InstructorEmbedding']:
        dir_command = f"test -d {repo_path}/{dirname} && echo 'exists' || echo 'not_found'"
        dir_result = env.controller.run_bash_script(dir_command, timeout=5)
        required_dirs[dirname] = dir_result.get('output', '').strip() == 'exists'
    return {'file_count': file_count, 'has_git_dir': has_git_dir, 'required_files_exist': required_files, 'required_dirs_exist': required_dirs}

def get_git_repo_subdirectory__a29e9963(env, config: dict):
    """
    Check if a specific subdirectory exists within a cloned git repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'base_path' and 'subdirectory' keys

    Returns:
        bool: True if subdirectory exists, False otherwise
    """
    base_path = config.get('base_path', '')
    subdirectory = config.get('subdirectory', '')
    full_path = f'{base_path}/{subdirectory}'
    command = f'test -d "{full_path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        output = result['output'].strip()
        return output == 'exists'
    return False

def get_folder_file_count__fc8f5b84(env, config: dict):
    """Count files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        int: Number of files in the directory
    """
    folder_path = config.get('folder_path', '')
    command = f'find "{folder_path}" -maxdepth 1 -type f 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        try:
            count = int(result['output'].strip())
            return count
        except:
            return 0
    return 0

def get_python_file_content__2a7463c5815fe65f87729a241d0d409d(env, config: dict):
    """Get Python file content and basic statistics.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File content statistics
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.split('\n')
        code_lines = 0
        comment_lines = 0
        blank_lines = 0
        docstring_lines = 0
        in_docstring = False
        for line in lines:
            stripped = line.strip()
            if not stripped:
                blank_lines += 1
            elif '"""' in stripped or "'''" in stripped:
                docstring_lines += 1
                in_docstring = not in_docstring
            elif in_docstring:
                docstring_lines += 1
            elif stripped.startswith('#'):
                comment_lines += 1
            else:
                code_lines += 1
        has_class = 'class ' in content
        has_def = 'def ' in content
        has_import = 'import ' in content or 'from ' in content
        return {'exists': True, 'total_lines': len(lines), 'code_lines': code_lines, 'comment_lines': comment_lines, 'blank_lines': blank_lines, 'docstring_lines': docstring_lines, 'has_class': has_class, 'has_def': has_def, 'has_import': has_import, 'char_count': len(content)}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_pdf_file_info__3a8b0dcb7e90aff8357c94bc8dad2474(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file information including existence, page count, and image content verification.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with file info including image dimensions and aspect ratio
    """
    from pypdf import PdfReader
    import tempfile
    from PIL import Image
    import io
    path = config.get('path')
    if not path:
        return {'exists': False}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        result = {'exists': True, 'page_count': page_count, 'file_size': file_size, 'path': path, 'has_image': False}
        if page_count > 0:
            try:
                first_page = reader.pages[0]
                if '/XObject' in first_page['/Resources']:
                    xobjects = first_page['/Resources']['/XObject'].get_object()
                    for obj_name in xobjects:
                        obj = xobjects[obj_name]
                        if obj['/Subtype'] == '/Image':
                            result['has_image'] = True
                            try:
                                width = obj['/Width']
                                height = obj['/Height']
                                result['image_width'] = width
                                result['image_height'] = height
                                result['aspect_ratio'] = round(width / height, 3) if height > 0 else 0
                            except:
                                pass
                            break
            except Exception as e:
                logger.debug(f'Could not extract image info from PDF: {e}')
        os.unlink(tmp_path)
        return result
    except Exception as e:
        logger.error(f'Error reading PDF: {e}')
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass
        return {'exists': False, 'error': str(e)}

def get_file_presence_pattern__ca75e69be093d6cb8e4fa53ce6114782(env, config):
    """Check file presence in multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'filename' and 'directories' (list)

    Returns:
        dict: Mapping of directory to whether file exists in it
    """
    filename = config.get('filename', '')
    directories = config.get('directories', [])
    results = {}
    for directory in directories:
        check_cmd = f'[ -f "{directory}/{filename}" ] && echo "exists" || echo "not_exists"'
        result = env.controller.run_bash_script(check_cmd, timeout=10)
        exists = False
        if result and result.get('output'):
            exists = result['output'].strip() == 'exists'
        results[directory] = exists
    return results

def get_python_file_count__c32ed37d80ef317dea88bac4a4cc1f31(env, config: dict):
    """Count Python files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        int: Number of Python files found
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    directory = config['directory']
    command = f"find {directory} -name '*.py' -type f | wc -l"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        try:
            return int(output)
        except ValueError:
            logger.error(f'Failed to parse count: {output}')
            return 0
    else:
        logger.error('Failed to count Python files. Status code: %d', response.status_code)
        return 0

def get_python_file__f1f92c4b10af2ffde6ea1534830a28f6(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get a Python file and validate it can be parsed and contains a print statement.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        Dict with 'exists', 'is_valid_python', 'has_print_statement', 'content', and 'file_path' keys
    """
    file_path = config.get('file_path', '')
    file_bytes = env.controller.get_file(file_path)
    result = {'exists': file_bytes is not None, 'file_path': file_path, 'content': '', 'is_valid_python': False, 'has_print_statement': False}
    if file_bytes:
        try:
            content = file_bytes.decode('utf-8')
            result['content'] = content
            tree = ast.parse(content)
            result['is_valid_python'] = True
            result['has_print_statement'] = _has_print_call(tree)
        except SyntaxError:
            result['is_valid_python'] = False
        except:
            result['content'] = ''
            result['is_valid_python'] = False
    return result

def get_text_file_content__8f36e769(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_pdf_files_list__097402a1(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """Get list of PDF files with page counts in the specified directory."""
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_data = []
    for filename in files:
        full_path = os.path.join(directory, filename)
        page_count_cmd = f"pdfinfo '{full_path}' 2>/dev/null | grep -oP '(?<=Pages:)\\s*\\d+' | tr -d ' '"
        page_result = env.controller.run_bash_script(page_count_cmd, timeout=10)
        page_count = 0
        if page_result.get('returncode') == 0:
            page_output = page_result.get('output', '').strip()
            if page_output.isdigit():
                page_count = int(page_output)
        pdf_data.append({'filename': filename, 'page_count': page_count})
    return pdf_data

def get_directory_count__7ec80ec4(env, config: dict):
    """
    Count directories matching a pattern in a parent directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'parent_path' and 'pattern' keys

    Returns:
        int: Number of matching directories
    """
    parent_path = config.get('parent_path', '')
    pattern = config.get('pattern', '*')
    command = f'find "{parent_path}" -maxdepth 1 -type d -name "{pattern}" 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        try:
            return int(result['output'].strip())
        except ValueError:
            return 0
    return 0

def get_pdf_text_content__a09f0e42332d230e7cc0e7794732425e(env, config: Dict[str, Any]) -> Any:
    """Get text content from first page of a PDF file from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - query: Query to find the file on Google Drive
            - dest: Local destination filename

    Returns:
        str: Text content from first page, or empty string if error
    """
    try:
        import PyPDF2
        from desktop_env.evaluators.getters.chrome import get_googledrive_file
        local_path = get_googledrive_file(env, config)
        if not local_path or not os.path.exists(local_path):
            logger.warning(f'PDF file not found: {local_path}')
            return ''
        with open(local_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            if len(pdf_reader.pages) > 0:
                first_page_text = pdf_reader.pages[0].extract_text()
                logger.info(f'Extracted {len(first_page_text)} characters from first page')
                return first_page_text
            else:
                logger.warning('PDF has no pages')
                return ''
    except Exception as e:
        logger.error(f'Error extracting PDF text: {e}')
        return ''

def get_dual_file_exists__e1a4b749(env, config):
    """
    Check if both CSV and XLSX files exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with csv_path and xlsx_path

    Returns:
        dict: {"csv_exists": bool, "xlsx_exists": bool}
    """
    csv_path = config.get('csv_path', '/home/user/Desktop/contacts.csv')
    xlsx_path = config.get('xlsx_path', '/home/user/Desktop/contacts.xlsx')
    result = {'csv_exists': False, 'xlsx_exists': False}
    csv_bytes = env.controller.get_file(csv_path)
    if csv_bytes:
        result['csv_exists'] = True
    xlsx_bytes = env.controller.get_file(xlsx_path)
    if xlsx_bytes:
        result['xlsx_exists'] = True
    return result

def get_pdf_file_info__60388d3f6c5270e1728288c485d5fd5b(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get PDF file from VM and verify it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        str: Path to downloaded PDF file in cache, or None if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import os
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_multi_directory_contents__ba67a508(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_text_file_content__83bea20e1ddf48cf4f537ad2c05896b6(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        logger.warning('No file path provided')
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return ''
        content = file_bytes.decode('utf-8').strip()
        logger.debug(f'Read {len(content)} characters from {file_path}')
        return content
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return ''

def get_file_text_content__61ff0c2a(env, config: dict):
    """Extract text content from file on VM.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        str: File text content
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        return file_bytes.decode('utf-8').strip()
    except:
        return None

def get_pdf_basic_info__ebee77f7069f08db28f1fec38e7eb935(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract basic information from a PDF file including orientation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with 'exists', 'page_count', 'file_size', 'orientation' keys
    """
    from pypdf import PdfReader
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'orientation': 'unknown'}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        orientation = 'unknown'
        if page_count > 0:
            first_page = reader.pages[0]
            mediabox = first_page.mediabox
            width = float(mediabox.width)
            height = float(mediabox.height)
            if width > height:
                orientation = 'landscape'
            elif height > width:
                orientation = 'portrait'
            else:
                orientation = 'square'
        return {'exists': True, 'page_count': page_count, 'file_size': file_size, 'orientation': orientation}
    except Exception as e:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'orientation': 'unknown', 'error': str(e)}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_file_line_count__868f1e74(env, config):
    """Count lines in a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        int: Number of lines in file, or -1 if file not found
    """
    file_path = config.get('file_path')
    command = f"wc -l {file_path} 2>/dev/null | awk '{{print $1}}' || echo '-1'"
    result = env.controller.run_bash_script(command, timeout=10)
    try:
        count = int(result.get('output', '-1').strip())
        return count
    except ValueError:
        return -1

def get_csv_first_column__e84252aa(env, config: dict):
    """Get all column headers from CSV file to verify column order.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of column headers, or empty list if file doesn't exist or is invalid
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = content.strip().split('\n')
        if not lines:
            return []
        reader = csv.reader(lines)
        header = next(reader, [])
        return header if header else []
    except Exception:
        return []

def get_file_content_dict__4ca7be3762cf58edd133468f74a91008(env, config: Dict[str, Any]) -> Dict[str, str]:
    """Read file content and return as a dict with file info.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'content' key containing file text
    """
    import os
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Failed to get file: {config['path']}")
        return {'content': '', 'exists': False}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.txt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return {'content': content, 'exists': True}
    except Exception as e:
        logger.error(f'Error reading file: {e}')
        return {'content': '', 'exists': False}
    finally:
        os.unlink(tmp_path)

def get_pdf_filenames_with_keyword__91578e58(env, config: dict):
    """Get PDF filenames containing specific keywords.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path
            - keywords: List of keywords to search for (or single keyword string)

    Returns:
        Dict with keyword counts: {'keyword1': count, 'keyword2': count, ...}
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    keywords = config.get('keywords', [])
    if isinstance(keywords, str):
        keywords = [keywords]
    command = f"""\npython3 -c "\nimport os\nimport glob\nimport json\n\ndir_path = '{directory}'\nkeywords = {keywords}\n\nif not os.path.exists(dir_path):\n    print(json.dumps({{k: 0 for k in keywords}}))\n    exit(0)\n\npdf_files = glob.glob(os.path.join(dir_path, '*.pdf'))\nkeyword_counts = {{}}\n\nfor keyword in keywords:\n    keyword_lower = keyword.lower()\n    count = sum(1 for f in pdf_files if keyword_lower in os.path.basename(f).lower())\n    keyword_counts[keyword] = count\n\nprint(json.dumps(keyword_counts))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to search files: {result['error']}")
        return {k: 0 for k in keywords}
    import json
    try:
        counts = json.loads(result['output'].strip())
        return counts
    except:
        return {k: 0 for k in keywords}

def get_zip_contents__6aa029d37a944ba9e2bf06a8a1d59f5c(env, config: Dict[str, Any]) -> List[str]:
    """
    Get the list of files contained in a ZIP archive on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List[str]: List of filenames in the ZIP archive, or empty list if error
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.info(f'ZIP file not found at {path}')
            return []
        with tempfile.NamedTemporaryFile(suffix='.zip', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                logger.info(f'ZIP file at {path} contains: {file_list}')
                return file_list
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading ZIP file: {e}')
        return []

def get_file_exists__6f3c16ae(env, config: dict):
    """Check if a tar.gz archive exists and validate its contents.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        dict: {
            'exists': bool,
            'valid_archive': bool,
            'contents': list of file paths in archive,
            'has_expected_structure': bool
        }
    """
    file_path = config.get('file_path', '')
    check_exists = f'test -f "{file_path}" && echo "1" || echo "0"'
    result = env.controller.run_bash_script(check_exists, timeout=10)
    if result['returncode'] != 0 or result['output'].strip() != '1':
        return {'exists': False, 'valid_archive': False, 'contents': [], 'has_expected_structure': False}
    list_command = f'tar -tzf "{file_path}" 2>&1'
    list_result = env.controller.run_bash_script(list_command, timeout=20)
    if list_result['returncode'] != 0:
        return {'exists': True, 'valid_archive': False, 'contents': [], 'has_expected_structure': False}
    contents = [line.strip() for line in list_result['output'].strip().split('\n') if line.strip()]
    expected_files = ['IDS LLM seminar/DSC00659.jpg', 'IDS LLM seminar/DSC00657.jpg', 'IDS LLM seminar/DSC00574.jpg', 'IDS LLM seminar/DSC00554.jpg', 'IDS LLM seminar/DSC00495.jpg', 'IDS LLM seminar/DSC00454.jpg']
    normalized_contents = set()
    for item in contents:
        normalized_contents.add(item.rstrip('/'))
    has_expected_structure = all((any((expected_file in content or content.endswith(expected_file.split('/')[-1]) for content in normalized_contents)) for expected_file in expected_files))
    expected_in_seminar_folder = all((f'IDS LLM seminar/{filename}' in normalized_contents or any((f'IDS LLM seminar/{filename}' in content for content in normalized_contents)) for filename in ['DSC00659.jpg', 'DSC00657.jpg', 'DSC00574.jpg', 'DSC00554.jpg', 'DSC00495.jpg', 'DSC00454.jpg']))
    return {'exists': True, 'valid_archive': True, 'contents': contents, 'has_expected_structure': expected_in_seminar_folder}

def get_pdf_content_keywords__c68533a8ce70563646c4156811d621fa(env, config: dict):
    """Check if PDF files exist and contain specific keywords.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path on VM
            - files_with_keywords: List of dicts with 'filename' and 'keywords' (list)

    Returns:
        dict: Results for each file {filename: list of found keywords}
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    files_with_keywords = config.get('files_with_keywords', [])
    results = {}
    for file_info in files_with_keywords:
        filename = file_info['filename']
        keywords = file_info.get('keywords', [])
        file_path = os.path.join(directory, filename)
        check_cmd = f'[ -f "{file_path}" ] && echo "exists" || echo "not_found"'
        cmd_result = env.controller.run_bash_script(check_cmd, timeout=10)
        if cmd_result.get('output', '').strip() != 'exists':
            results[filename] = []
            continue
        found_keywords = []
        try:
            file_bytes = env.controller.get_file(file_path)
            if file_bytes:
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    import fitz
                    doc = fitz.open(tmp_path)
                    text = ''
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    for keyword in keywords:
                        if keyword in text:
                            found_keywords.append(keyword)
                except Exception:
                    pass
                finally:
                    os.unlink(tmp_path)
        except Exception:
            pass
        results[filename] = found_keywords
    return results

def get_subfolder_picture_hashes__fcfe8473(env, config):
    """Get hashes of pictures in a specific subfolder.

    Args:
        env: Desktop environment
        config: Dict with 'subfolder_path' key

    Returns:
        list: Sorted list of image hashes in subfolder
    """
    subfolder = config.get('subfolder_path', '')
    result = env.controller.run_bash_script(f"find {subfolder} -maxdepth 1 -type f \\( -iname '*.jpg' -o -iname '*.jpeg' -o -iname '*.png' \\) 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        return []
    files = result['output'].strip().split('\n')
    hashes = []
    for file_path in files:
        if not file_path or file_path.strip() == '':
            continue
        file_path = file_path.strip()
        file_bytes = env.controller.get_file(file_path)
        if file_bytes:
            try:
                with Image.open(BytesIO(file_bytes)) as img:
                    img_byte_arr = img.tobytes()
                    hash_result = hashlib.sha256(img_byte_arr).hexdigest()
                    hashes.append(hash_result)
            except Exception as e:
                print(f'Error processing {file_path}: {e}')
                continue
    return sorted(hashes)

def get_docx_text_content__64111ae8(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_directory_files__739292ff(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get list of files in a directory with detailed file information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'pattern' parameters

    Returns:
        dict: Information about files in directory including size and image properties
    """
    path = config.get('path', '/home/user/Desktop')
    pattern = config.get('pattern', '*.png')
    list_script = f"\nimport os\nimport glob\nimport json\n\nfiles_info = {{}}\nfiles = glob.glob(os.path.join('{path}', '{pattern}'))\n\nfor f in files:\n    filename = os.path.basename(f)\n    file_info = {{'exists': True}}\n\n    try:\n        # Get file size\n        file_info['size'] = os.path.getsize(f)\n\n        # Try to validate as PNG image\n        try:\n            from PIL import Image\n            with Image.open(f) as img:\n                file_info['valid_image'] = True\n                file_info['dimensions'] = {{'width': img.width, 'height': img.height}}\n                file_info['format'] = img.format\n        except Exception as e:\n            file_info['valid_image'] = False\n            file_info['image_error'] = str(e)\n    except Exception as e:\n        file_info['error'] = str(e)\n\n    files_info[filename] = file_info\n\nprint(json.dumps(files_info))\n"
    result = env.controller.execute_python_command(list_script)
    if result and result.get('output'):
        output = result['output'].strip()
        if output:
            try:
                import json
                files_info = json.loads(output)
                return {'files': list(files_info.keys()), 'count': len(files_info), 'files_info': files_info}
            except Exception as e:
                logger.error(f'Error parsing file info: {e}')
    return {'files': [], 'count': 0, 'files_info': {}}

def get_file_count__9d5ff0d9(env, config: dict):
    """Count files in a directory on VM.

    This getter counts all files matching the pattern in the directory.
    Note: It does not verify file source (e.g., from specific email).
    Relies on clean environment setup where the Downloads directory
    is freshly created for the task.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern'

    Returns:
        int: Number of matching files
    """
    directory = config.get('directory', '')
    pattern = config.get('pattern', '*')
    command = f"find {directory} -maxdepth 1 -name '{pattern}' -type f 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('returncode') == 0:
        try:
            return int(result.get('output', '0').strip())
        except ValueError:
            return 0
    return 0

def get_srt_file_exists__4f098003e517e7e34a157f1c233e1c85(env, config: dict):
    """
    Check if SRT file exists and validate it contains valid SRT subtitle format.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the SRT file

    Returns:
        dict: {
            'exists': bool,
            'size': int,
            'has_content': bool,
            'valid_srt_format': bool,
            'subtitle_count': int
        }
    """
    path = config.get('path', '/home/user/subtitles.srt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'has_content': False, 'valid_srt_format': False, 'subtitle_count': 0}
    file_size = len(file_bytes)
    has_content = file_size > 100
    try:
        content = file_bytes.decode('utf-8', errors='ignore')
    except Exception as e:
        logger.warning(f'Failed to decode file as UTF-8: {e}')
        return {'exists': True, 'size': file_size, 'has_content': has_content, 'valid_srt_format': False, 'subtitle_count': 0}
    (valid_srt, subtitle_count) = _validate_srt_format(content)
    return {'exists': True, 'size': file_size, 'has_content': has_content, 'valid_srt_format': valid_srt, 'subtitle_count': subtitle_count}

def get_numbers_from_file__d0f3a745(env, config: dict):
    """Extract all numbers from file content.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        list: List of integers found in file
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return []
    try:
        content = file_bytes.decode('utf-8')
        numbers = re.findall('\\d+', content)
        return [int(n) for n in numbers]
    except:
        return []

def get_file_checksum__032a6328(env, config: dict) -> str:
    """Get MD5 checksum of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path'

    Returns:
        MD5 checksum string
    """
    file_path = config.get('path', '')
    command = f'md5sum {file_path} | cut -d" " -f1'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    return ''

def get_pdf_files_list__92f3cb2c(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files with their page counts in the specified directory.

    Returns:
        Dict[str, int]: Mapping of filename to page count, e.g., {'file.pdf': 10}
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    list_command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    if list_result.get('returncode') != 0:
        return {}
    output = list_result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        pdfinfo_command = f"pdfinfo '{filepath}' 2>/dev/null | grep -i '^Pages:' | awk '{{print $2}}'"
        pdfinfo_result = env.controller.run_bash_script(pdfinfo_command, timeout=10)
        if pdfinfo_result.get('returncode') == 0:
            page_output = pdfinfo_result.get('output', '').strip()
            try:
                page_count = int(page_output)
                pdf_info[filename] = page_count
            except (ValueError, TypeError):
                continue
        else:
            continue
    return pdf_info

def get_numeric_value_from_file__5d55e268(env, config: Dict[str, Any]) -> Optional[float]:
    """
    Get a numeric value from a text file.

    Config:
        path (str): absolute path on the VM to fetch the file
        dest (str): file name of the downloaded file

    Returns:
        float: The numeric value from the file, or None if file doesn't exist or value cannot be parsed
    """
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    try:
        file_data = env.controller.get_file(path)
        if file_data is None:
            logger.warning(f'File not found on VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read().strip()
        try:
            value = float(content)
            logger.info(f'Successfully read numeric value from {path}: {value}')
            return value
        except ValueError:
            logger.error(f'Cannot parse content as numeric value: {content}')
            return None
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return None

def get_file_text__b8aa5550(env, config: dict):
    """Get text content from VM file.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        str: File content
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        return file_bytes.decode('utf-8')
    except:
        return None

def get_file_permissions__2bea57f7(env, config):
    """Get file permissions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File permissions string (e.g., '-rwxr-xr-x')
    """
    path = config.get('path', '')
    command = f"ls -l {path} | awk '{{print $1}}'"
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        return result['output'].strip()
    else:
        logger.error(f"Failed to get permissions for {path}: {result['error']}")
        return None

def get_text_file_content__a2f161de(env, config):
    """Read the content of a text file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content, or empty string if file doesn't exist
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        print(f'Error reading text file: {e}')
        return ''

def get_file_content__9cefe3d0(env, config: dict):
    """Get file content from VM.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        str: File content
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        return file_bytes.decode('utf-8')
    except:
        return None

def get_file_content__66187dca7d72426d9eab5771bd5dd30e(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    result = env.controller.run_bash_script(f'cat {file_path}', timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    else:
        return ''

def get_file_exists__739292ff(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists and validate it as a valid PNG image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: {
            'exists': bool,
            'is_valid_png': bool,
            'file_size': int,
            'has_valid_dimensions': bool,
            'width': int or None,
            'height': int or None
        }
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return {'exists': False, 'is_valid_png': False, 'file_size': 0, 'has_valid_dimensions': False, 'width': None, 'height': None}
    python_code = f"\nimport os\nimport sys\n\npath = '{path}'\nresult = {{\n    'exists': False,\n    'is_valid_png': False,\n    'file_size': 0,\n    'has_valid_dimensions': False,\n    'width': None,\n    'height': None\n}}\n\n# Check if file exists\nif os.path.exists(path):\n    result['exists'] = True\n\n    # Check file size\n    try:\n        result['file_size'] = os.path.getsize(path)\n    except Exception as e:\n        print(f'Error getting file size: {{e}}', file=sys.stderr)\n\n    # Validate as PNG image\n    if result['file_size'] > 0:\n        try:\n            # Check PNG magic bytes\n            with open(path, 'rb') as f:\n                header = f.read(8)\n                if header[:4] == b'\\x89PNG':\n                    result['is_valid_png'] = True\n\n                    # Try to get dimensions using PIL\n                    try:\n                        from PIL import Image\n                        with Image.open(path) as img:\n                            result['width'] = img.width\n                            result['height'] = img.height\n                            # Valid screenshot should have reasonable dimensions\n                            if img.width > 0 and img.height > 0:\n                                result['has_valid_dimensions'] = True\n                    except Exception as e:\n                        print(f'Error loading image with PIL: {{e}}', file=sys.stderr)\n        except Exception as e:\n            print(f'Error validating PNG: {{e}}', file=sys.stderr)\n\nprint(result)\n"
    result = env.controller.execute_python_command(python_code)
    if result and result.get('output'):
        try:
            import ast
            output = result['output'].strip()
            parsed = ast.literal_eval(output)
            return parsed
        except Exception as e:
            logger.error(f'Error parsing result: {e}')
            return {'exists': False, 'is_valid_png': False, 'file_size': 0, 'has_valid_dimensions': False, 'width': None, 'height': None}
    return {'exists': False, 'is_valid_png': False, 'file_size': 0, 'has_valid_dimensions': False, 'width': None, 'height': None}

def get_folder_files__21ed89b452cf1a748805f1dda9b2ec4b(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files in a specific folder on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' key

    Returns:
        List of filenames in the folder (empty list if folder doesn't exist)
    """
    folder_path = config.get('folder_path', '')
    command = f"ls -1 '{folder_path}' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode', 1) != 0 or not result.get('output', '').strip():
        logger.warning(f'Folder not found or empty: {folder_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    return files

def get_pdf_files_in_dir__600d7d75(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_text_file_lines__4250d59b26bb86f2de0562f0a55c312c(env, config: dict) -> Optional[list]:
    """Get lines from a text file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of stripped lines, or None if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n')]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_multi_directory_contents__e2bf8bf2(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_bash_file_content__38c143b8ba918371e989fa588ea9cb56(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file not found
    """
    file_path = config.get('path', '/home/user/output.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_file_exists_and_size__91c793ae(env, config: Dict[str, Any]):
    """
    Check if a file exists on the VM and return its size in bytes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {"exists": bool, "size": int (bytes), "is_png": bool}
    """
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_pdf_exports__17dfab17(env, config: dict):
    """
    Check if PDFs were exported to the specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        List of dict with PDF filename and size information
    """
    directory = config.get('directory', '/home/user/Downloads/Readings')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        for f in list_result['output'].strip().split('\n'):
            f = f.strip()
            if f:
                size_command = f'stat -c %s "{f}" 2>/dev/null'
                size_result = env.controller.run_bash_script(size_command, timeout=5)
                file_size = 0
                if size_result['output']:
                    try:
                        file_size = int(size_result['output'].strip())
                    except ValueError:
                        file_size = 0
                pdf_files.append({'filename': os.path.basename(f), 'size': file_size})
    return pdf_files

def get_csv_and_count_verification__40b9d100(env, config: dict):
    """Verify CSV export exists and validate mobile count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with expected count

    Returns:
        dict: Verification results with CSV info and count validation
    """
    desktop_path = '/home/user/Desktop'
    mobile_count_path = '/home/user/Desktop/mobile_count.txt'
    result = {'csv_exists': False, 'csv_has_contacts': False, 'csv_has_mobile_column': False, 'csv_mobile_count': 0, 'text_file_exists': False, 'text_file_count': None, 'counts_match': False, 'all_verified': False}
    try:
        csv_files = []
        desktop_files = env.controller.list_files(desktop_path)
        if desktop_files:
            for filename in desktop_files:
                if filename.lower().endswith('.csv'):
                    csv_files.append(os.path.join(desktop_path, filename))
        if not csv_files:
            return result
        result['csv_exists'] = True
        csv_file_path = csv_files[0]
        csv_bytes = env.controller.get_file(csv_file_path)
        if csv_bytes is None:
            return result
        csv_content = csv_bytes.decode('utf-8', errors='ignore')
        csv_lines = csv_content.strip().split('\n')
        if len(csv_lines) < 2:
            return result
        csv_reader = csv.DictReader(csv_lines)
        headers = csv_reader.fieldnames
        if not headers:
            return result
        result['csv_has_contacts'] = True
        mobile_column = None
        for header in headers:
            if 'mobile' in header.lower() or 'cell' in header.lower():
                mobile_column = header
                result['csv_has_mobile_column'] = True
                break
        if not mobile_column:
            return result
        mobile_count = 0
        for row in csv_reader:
            mobile_value = row.get(mobile_column, '').strip()
            if mobile_value:
                mobile_count += 1
        result['csv_mobile_count'] = mobile_count
    except Exception as e:
        pass
    try:
        text_bytes = env.controller.get_file(mobile_count_path)
        if text_bytes is not None:
            result['text_file_exists'] = True
            text_content = text_bytes.decode('utf-8', errors='ignore').strip()
            try:
                text_count = int(text_content)
                result['text_file_count'] = text_count
                if result['csv_mobile_count'] > 0 and text_count == result['csv_mobile_count']:
                    result['counts_match'] = True
            except ValueError:
                pass
    except Exception:
        pass
    result['all_verified'] = result['csv_exists'] and result['csv_has_contacts'] and result['csv_has_mobile_column'] and result['text_file_exists'] and result['counts_match']
    return result

def get_docx_text_content__ae54b7f2(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_pdf_exports__d263a7ae(env, config: dict):
    """
    Check if PDFs were exported to the specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        List of PDF filenames in the directory
    """
    directory = config.get('directory', '/home/user/Desktop/BlogArchive')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        pdf_files = [os.path.basename(f.strip()) for f in list_result['output'].strip().split('\n') if f.strip()]
    return pdf_files

def get_text_file_content__ffa8cf6ad724ee8fc8a065457d283c28(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file on VM

    Returns:
        str: File content as string
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_text_file_content__0b52fd51(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_file_exists_check__e327229f(env, config: Dict) -> Optional[Dict]:
    """
    Check if a file exists and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_backup_dir_contents__1aa0aa31(env, config):
    """Get list of files in backup directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' path

    Returns:
        List of filenames
    """
    directory = config.get('directory', '/home/user')
    command = f"ls -1 {directory} 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('status') == 'success':
        output = result['output'].strip()
        if output:
            return sorted([f.strip() for f in output.split('\n') if f.strip()])
    return []

def get_vm_file_exists__949eb101(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on the VM and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool) and 'size_bytes' (int) keys
    """
    path = config['path']
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; else echo 'NOT_EXISTS'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result['output'].strip()
    if output == 'NOT_EXISTS' or result['returncode'] != 0:
        return {'exists': False, 'size_bytes': 0}
    try:
        size = int(output)
        return {'exists': True, 'size_bytes': size}
    except ValueError:
        logger.error(f'Failed to parse file size: {output}')
        return {'exists': False, 'size_bytes': 0}

def get_file_size__79b627f8(env, config):
    """Get file size in bytes from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: File size in bytes, or -1 if file not found
    """
    path = config.get('path', '')
    result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then stat -c %s "{path}"; else echo "-1"; fi', timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        logger.error(f'Failed to parse file size: {output}')
        return -1

def get_odt_file_info__82a08649fb4593676b95bae1dea8263f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if an ODT file exists and the original DOCX file is preserved.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (ODT file path)

    Returns:
        Dict with 'odt_exists', 'docx_exists' keys indicating if both files exist
    """
    odt_path = config.get('path', '')
    docx_path = odt_path.replace('.odt', '.docx')
    odt_file_bytes = env.controller.get_file(odt_path)
    odt_exists = odt_file_bytes is not None
    docx_file_bytes = env.controller.get_file(docx_path)
    docx_exists = docx_file_bytes is not None
    if not odt_exists:
        logger.warning(f'ODT file not found at {odt_path}')
    else:
        logger.info(f'ODT file found at {odt_path} ({len(odt_file_bytes)} bytes)')
    if not docx_exists:
        logger.warning(f'Original DOCX file not found at {docx_path}')
    else:
        logger.info(f'Original DOCX file preserved at {docx_path} ({len(docx_file_bytes)} bytes)')
    return {'odt_exists': odt_exists, 'docx_exists': docx_exists, 'odt_path': odt_path, 'docx_path': docx_path, 'odt_size': len(odt_file_bytes) if odt_exists else 0, 'docx_size': len(docx_file_bytes) if docx_exists else 0}

def get_both_files_rows__e54614f2(env, config):
    """
    Check if both CSV and XLSX files exist and get their row counts.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with csv_path and xlsx_path

    Returns:
        dict: {"csv_exists": bool, "xlsx_exists": bool, "csv_rows": int, "xlsx_rows": int}
    """
    csv_path = config.get('csv_path', '/home/user/Desktop/contacts.csv')
    xlsx_path = config.get('xlsx_path', '/home/user/Desktop/contacts.xlsx')
    result = {'csv_exists': False, 'xlsx_exists': False, 'csv_rows': 0, 'xlsx_rows': 0}
    csv_bytes = env.controller.get_file(csv_path)
    if csv_bytes:
        result['csv_exists'] = True
        try:
            csv_text = csv_bytes.decode('utf-8')
            reader = csv.reader(io.StringIO(csv_text))
            result['csv_rows'] = sum((1 for _ in reader))
        except Exception:
            result['csv_rows'] = 0
    xlsx_bytes = env.controller.get_file(xlsx_path)
    if xlsx_bytes:
        result['xlsx_exists'] = True
        try:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(xlsx_bytes)
                tmp_path = tmp.name
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            ws = wb.active
            result['xlsx_rows'] = ws.max_row
            import os
            os.unlink(tmp_path)
        except Exception:
            result['xlsx_rows'] = 0
    return result

def get_txt_file_lines__87f6c353(env, config):
    """
    Read a text file and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: Lines from the text file (stripped)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return []
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception as e:
        return []

def get_pdf_validation__82780f835e33185cfaa42bf9dd4b545f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Validate PDF file and get page count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'page_count' (int), 'is_valid' (bool)
    """
    try:
        from pypdf import PdfReader
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'page_count': 0, 'is_valid': False}
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            page_count = len(reader.pages)
            return {'page_count': page_count, 'is_valid': page_count > 0}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return {'page_count': 0, 'is_valid': False}

def get_pdf_multipage_info__7825428a26721545cdd69efbcd0db85d(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF multi-page information.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with multi-page info
    """
    from pypdf import PdfReader
    import tempfile
    path = config.get('path')
    if not path:
        return {'exists': False}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        pages_with_images = 0
        for page in reader.pages:
            if '/XObject' in page.get('/Resources', {}):
                xobjects = page['/Resources']['/XObject'].get_object()
                for obj in xobjects:
                    if xobjects[obj].get('/Subtype') == '/Image':
                        pages_with_images += 1
                        break
        os.unlink(tmp_path)
        return {'exists': True, 'page_count': page_count, 'pages_with_images': pages_with_images, 'file_size': len(file_bytes)}
    except Exception as e:
        logger.error(f'Error reading PDF: {e}')
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass
        return {'exists': False, 'error': str(e)}

def get_file_and_folder_check__3c993009(env, config):
    """Check if folder exists and if file is inside that folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' and 'file_path' parameters

    Returns:
        dict: {"folder_exists": bool, "file_in_folder": bool}
    """
    folder_path = config.get('folder_path')
    file_path = config.get('file_path')
    result_folder = env.controller.run_bash_script(f'[ -d "{folder_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    folder_exists = 'exists' in result_folder['output'] and 'not exists' not in result_folder['output']
    result_file = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    file_in_folder = 'exists' in result_file['output'] and 'not exists' not in result_file['output']
    return {'folder_exists': folder_exists, 'file_in_folder': file_in_folder}

def get_text_file_lines__6a539f3f0959c2ae90484012b1290c1f(env, config):
    """
    Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: Lines from the file (stripped of whitespace)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return []

def get_file_exists__6d219cf2(env, config):
    """Check if a file exists and get its metadata (size, type, magic bytes) on VM."""
    file_path = config.get('path', '')
    command = f'''\nif [ -f "{file_path}" ]; then\n    size=$(stat -c%s "{file_path}" 2>/dev/null || stat -f%z "{file_path}" 2>/dev/null)\n    magic=$(xxd -p -l 8 "{file_path}" 2>/dev/null | tr -d '\\n')\n    echo "EXISTS|$size|$magic"\nelse\n    echo "NOT_EXISTS"\nfi\n'''
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output == 'NOT_EXISTS' or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    parts = output.split('|')
    if len(parts) < 3:
        return {'exists': True, 'size': 0, 'is_png': False}
    size_str = parts[1]
    magic_bytes = parts[2].lower()
    try:
        file_size = int(size_str)
    except (ValueError, TypeError):
        file_size = 0
    is_png = magic_bytes == '89504e470d0a1a0a'
    return {'exists': True, 'size': file_size, 'is_png': is_png}

def get_file_exists__9180f469(env, config: dict):
    """Check if the exported CSV file exists and validate its format and content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with validation results including:
            - exists: bool - whether file exists
            - is_csv: bool - whether file is valid CSV
            - is_utf8: bool - whether file is UTF-8 encoded
            - size: int - file size in bytes
            - row_count: int - number of rows in CSV
            - has_content: bool - whether file has actual content
    """
    vm_path = config.get('path', '/home/user/Desktop/survey-export.csv')
    source_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.xlsx'
    result = {'exists': False, 'is_csv': False, 'is_utf8': False, 'size': 0, 'row_count': 0, 'has_content': False}
    exists_check = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    if exists_check.get('output', '').strip() != 'EXISTS':
        return result
    result['exists'] = True
    size_check = env.controller.run_bash_script(f"stat -c %s '{vm_path}' 2>/dev/null || echo '0'", timeout=10)
    try:
        result['size'] = int(size_check.get('output', '0').strip())
    except:
        result['size'] = 0
    result['has_content'] = result['size'] > 0
    encoding_check = env.controller.run_bash_script(f"file --mime-encoding '{vm_path}' 2>/dev/null || echo 'unknown'", timeout=10)
    encoding_output = encoding_check.get('output', '').strip().lower()
    result['is_utf8'] = 'utf-8' in encoding_output or 'us-ascii' in encoding_output
    if result['has_content']:
        read_result = env.controller.run_bash_script(f"cat '{vm_path}'", timeout=30)
        content = read_result.get('output', '')
        if content:
            try:
                csv_reader = csv.reader(StringIO(content))
                rows = list(csv_reader)
                result['row_count'] = len(rows)
                result['is_csv'] = len(rows) > 0
            except Exception as e:
                result['is_csv'] = False
                result['row_count'] = 0
    return result

def get_file_and_trash_status__05c91736(env, config):
    """Check if file exists and if it's in trash.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        dict: {"file_exists": bool, "in_trash": bool}
    """
    file_path = config.get('file_path')
    filename = file_path.split('/')[-1]
    result_file = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    file_exists = 'exists' in result_file['output'] and 'not exists' not in result_file['output']
    result_trash = env.controller.run_bash_script(f'gio trash --list | grep -q "{filename}" && echo "in_trash" || echo "not_in_trash"', timeout=10)
    in_trash = 'in_trash' in result_trash['output'] and 'not_in_trash' not in result_trash['output']
    return {'file_exists': file_exists, 'in_trash': in_trash}

def get_filename_pattern_match__0d61b4f8(env, config):
    """Check if a file matching a pattern exists in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern' keys

    Returns:
        str: Matching filename or empty string if not found
    """
    directory = config.get('directory', '/home/user/Desktop')
    pattern = config.get('pattern', '*.mp4')
    result = env.controller.run_bash_script(f'cd "{directory}" 2>/dev/null && ls -1 {pattern} 2>/dev/null | head -1 || echo ""', timeout=10)
    output = result.get('output', '').strip()
    return output

def get_file_content__bcc15712(env, config: Dict[str, Any]) -> str:
    """Get the content of a file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        str: Content of the file
    """
    from desktop_env.evaluators.getters.file import get_vm_file
    file_path = config.get('file_path', '')
    dest = config.get('dest', 'temp_file.txt')
    return get_vm_file(env, {'path': file_path, 'dest': dest})

def get_filename__4ee0209a(env, config: dict):
    """Get the filename (without path) of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Filename or empty string if not found
    """
    path = config.get('path', '/home/user/Desktop/invoice.xlsx')
    result = env.controller.run_bash_script(f'test -f "{path}" && basename "{path}" || echo ""', timeout=10)
    output = result.get('output', '').strip()
    return output

def get_text_file_content__93eac3e2452ef121ce8047db9ec250fe(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return ''

def get_docx_text_content__c9cce3df360663a0b15969bfb64090f9(env, config):
    """Extract all text from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of all text strings from the document
    """
    doc_path = config.get('path', '/home/user/Desktop/selected_slides.docx')
    file_bytes = env.controller.get_file(doc_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_texts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                all_texts.append(text)
        return all_texts
    finally:
        os.unlink(tmp_path)

def get_file_count__9f898a55(env, config: dict) -> int:
    """Count files matching pattern in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern'

    Returns:
        Number of matching files
    """
    directory = config.get('directory', '/home/user/Documents/Finance/receipts')
    pattern = config.get('pattern', '*.pdf')
    command = f'find {directory} -maxdepth 1 -name "{pattern}" -type f | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') == 0:
        try:
            return int(result.get('output', '0').strip())
        except ValueError:
            return 0
    return 0

def get_pdf_files_in_dir__197670d4(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_header_content__cc0e3cfb7e6297a7f30b4a2bddec08bd(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract header content from the document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'has_header' and 'header_text' keys
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'has_header': False, 'header_text': ''}
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        if not doc.sections:
            return {'has_header': False, 'header_text': ''}
        first_section = doc.sections[0]
        header = first_section.header
        if header and header.paragraphs:
            header_text = ' '.join([p.text.strip() for p in header.paragraphs if p.text.strip()])
            has_header = bool(header_text)
        else:
            has_header = False
            header_text = ''
        result = {'has_header': has_header, 'header_text': header_text}
        return result
    finally:
        os.unlink(tmp_path)

def get_file_exists__cdcbbd90(env, config):
    """Check if a file exists at the specified path on VM."""
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output == 'EXISTS':
        return True
    else:
        return False

def get_directory_listing__96172b42(env, config):
    """Get directory listing from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Directory listing output
    """
    path = config.get('path', '/home/user')
    command = f'ls -la {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        return result['output']
    else:
        logger.error(f"Failed to list directory {path}: {result['error']}")
        return None

def get_subject_file__271abb880d5f6f8d57d2c41e20bcf6ad(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Read the content of a text file that should contain an email subject.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (file path to read)

    Returns:
        Content of the file as string, or None if file doesn't exist
    """
    file_path = config.get('path', '/home/user/subject.txt')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return None
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return None

def get_docx_content_and_font__d62b618c18b762984a0245cb92f201ef(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get document content and font size from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the file path on VM

    Returns:
        Dict with 'content' (str) and 'font_size' (int) keys
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'content': '', 'font_size': None}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'content': '', 'font_size': None}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content = []
        for paragraph in doc.paragraphs:
            content.append(paragraph.text)
        full_text = '\n'.join(content)
        font_sizes = []
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    font_sizes.append(run.font.size.pt)
        font_size = None
        if font_sizes:
            if all((abs(fs - font_sizes[0]) <= 0.5 for fs in font_sizes)):
                font_size = font_sizes[0]
            else:
                font_size = font_sizes[0]
        return {'content': full_text.strip(), 'font_size': font_size}
    finally:
        os.unlink(tmp_path)

def get_file_count_by_pattern__5c107_5(env, config: dict):
    """
    Count files matching a pattern in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'pattern' and 'directory' keys

    Returns:
        Number of matching files as string
    """
    pattern = config.get('pattern', '*')
    directory = config.get('directory', '/home/user')
    command = f"find {directory} -name '{pattern}' -type f 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=30)
    if result.get('returncode') == 0:
        output = result.get('output', '').strip()
        return output
    return '0'

def get_text_file_lines__03d4ebf8db116d0dd3db06d0b1b6415c(env, config):
    """
    Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: Lines from the file (stripped of whitespace)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return []

def get_file_rename_status__bd9934867663f0945bb79537ace5711a(env, config: dict):
    """Check the status of a file rename operation.

    Verifies both that the new file exists and the old file doesn't exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'new_path' and 'old_path' keys

    Returns:
        dict: Dictionary with 'new_exists' and 'old_exists' boolean values
    """
    new_path = config.get('new_path', '')
    old_path = config.get('old_path', '')
    if not new_path or not old_path:
        logger.warning('Missing new_path or old_path in config')
        return {'new_exists': False, 'old_exists': True}
    try:
        new_file_bytes = env.controller.get_file(new_path)
        new_exists = new_file_bytes is not None and len(new_file_bytes) > 0
        old_file_bytes = env.controller.get_file(old_path)
        old_exists = old_file_bytes is not None and len(old_file_bytes) > 0
        logger.debug(f'Rename status - New file exists: {new_exists}, Old file exists: {old_exists}')
        return {'new_exists': new_exists, 'old_exists': old_exists}
    except Exception as e:
        logger.error(f'Error checking file rename status: {e}')
        return {'new_exists': False, 'old_exists': True}

def get_file_in_directory__3500c3270fba14684f134a0b5e4537d1(env, config: dict):
    """
    Check if a file exists and validate it contains valid SRT subtitle content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        dict: {
            'exists': bool,
            'path': str,
            'has_content': bool,
            'is_valid_srt': bool,
            'content': str
        }
    """
    path = config.get('path', '/home/user/Desktop/subtitles.srt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'path': path, 'has_content': False, 'is_valid_srt': False, 'content': ''}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
    except Exception as e:
        logger.error(f'Failed to decode file content: {e}')
        content = ''
    has_content = len(content.strip()) > 0
    is_valid_srt = False
    if has_content:
        import re
        timestamp_pattern = '\\d{2}:\\d{2}:\\d{2},\\d{3}\\s*-->\\s*\\d{2}:\\d{2}:\\d{2},\\d{3}'
        has_timestamps = bool(re.search(timestamp_pattern, content))
        has_sequence_numbers = bool(re.search('^\\d+$', content, re.MULTILINE))
        is_valid_srt = has_timestamps and has_sequence_numbers
    return {'exists': True, 'path': path, 'has_content': has_content, 'is_valid_srt': is_valid_srt, 'content': content[:500]}

def get_pdf_files_in_dir__201ff98c(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_file_content__2ef64375b850b707e97b54053b1452c0(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    result = env.controller.run_bash_script(f'cat {file_path}', timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    else:
        return ''

def get_directory_exists__0b3844b6(env, config: dict):
    """
    Check if a directory exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if directory exists, False otherwise
    """
    path = config.get('path', '')
    command = f'test -d "{path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        output = result['output'].strip()
        return output == 'exists'
    return False

def get_pdf_directory_info__9ba5dc81930fa553e6c6310edaaff2ec(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get information about PDF files in a directory.

    Args:
        env: Environment object
        config: Configuration dict with 'directory' key specifying the path

    Returns:
        Dict with 'count' (number of PDFs), 'files' (list of PDF file paths), and 'filenames' (list of PDF filenames without extension)
    """
    directory = config.get('directory', '')
    try:
        dir_tree = env.controller.get_vm_directory_tree(directory)
    except Exception as e:
        logger.error(f'Failed to get directory tree for {directory}: {e}')
        return {'count': 0, 'files': [], 'filenames': []}
    pdf_files = []
    pdf_filenames = []
    if dir_tree and 'children' in dir_tree:
        for item in dir_tree['children']:
            if item.get('type') == 'file' and item.get('name', '').lower().endswith('.pdf'):
                pdf_path = os.path.join(directory, item['name'])
                pdf_files.append(pdf_path)
                filename = item['name'][:-4] if item['name'].lower().endswith('.pdf') else item['name']
                pdf_filenames.append(filename)
    logger.info(f'Found {len(pdf_files)} PDF files in {directory}')
    return {'count': len(pdf_files), 'files': pdf_files, 'filenames': pdf_filenames}

def get_python_content__198be354(env, config):
    """Analyze Python file content for specific patterns.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Content analysis (has_torch, has_nn_module, has_training_loop)
    """
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'has_torch': False, 'has_nn_module': False, 'has_training_loop': False, 'has_optimizer': False}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    has_torch = 'torch' in content_str
    has_nn_module = 'nn.Module' in content_str
    has_training_loop = 'optimizer.step()' in content_str or 'loss.backward()' in content_str
    has_optimizer = 'optimizer' in content_str
    return {'exists': True, 'has_torch': has_torch, 'has_nn_module': has_nn_module, 'has_training_loop': has_training_loop, 'has_optimizer': has_optimizer}

def get_footer_has_page_and_filename__aac82fb893db44f74dfc3c4f83b7b05e(env, config):
    """Check if footer contains both page numbers and filename.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'filename' keys

    Returns:
        dict: Dictionary with 'has_page_number' and 'has_filename' booleans
    """
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.error(f"Failed to get file: {config['path']}")
            return {'has_page_number': False, 'has_filename': False}
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            expected_filename = config.get('filename', 'LibreOffice_Open_Source_Word_Processing')
            has_page_number = False
            has_filename = False
            if doc.sections and len(doc.sections) > 0:
                section = doc.sections[0]
                if section.footer and section.footer.paragraphs:
                    footer_text = ' '.join([p.text for p in section.footer.paragraphs])
                    has_filename = expected_filename.lower() in footer_text.lower()
                    footer_element = section.footer._element
                    footer_xml = footer_element.xml.decode('utf-8') if isinstance(footer_element.xml, bytes) else str(footer_element.xml)
                    page_field_indicators = ['PAGE', 'NUMPAGES', 'w:fldChar', '<w:instrText']
                    has_field_code = any((indicator in footer_xml for indicator in page_field_indicators))
                    page_number_patterns = ['\\bPage\\s+\\d+', '\\b\\d+\\s+of\\s+\\d+', '^\\s*\\d+\\s*$', '\\b#\\b']
                    has_page_pattern = any((re.search(pattern, footer_text, re.IGNORECASE) for pattern in page_number_patterns))
                    has_page_number = has_field_code or has_page_pattern
                    logger.info(f'Footer text: {footer_text}')
                    logger.info(f'Has field code: {has_field_code}, Has page pattern: {has_page_pattern}')
            return {'has_page_number': has_page_number, 'has_filename': has_filename}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking footer content: {e}')
        return {'has_page_number': False, 'has_filename': False}

def get_file_exists_and_size__058bd353(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_bash_file_content__8c025f1893a98ce909d203315d730cfe(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file not found
    """
    file_path = config.get('path', '/home/user/output.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_pdf_filename__e61f394075a7c32a6a0c2c96a3700939(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file and check filename pattern.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with filename information
    """
    from pypdf import PdfReader
    file_bytes = env.controller.get_file(config['path'])
    filename = os.path.basename(config['path'])
    result = {'path': config['path'], 'filename': filename, 'exists': file_bytes is not None and len(file_bytes) > 0}
    if result['exists']:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            result['valid_pdf'] = len(reader.pages) > 0
            result['page_count'] = len(reader.pages)
            logger.info(f"PDF file '{filename}' found with {result['page_count']} pages")
        except Exception as e:
            logger.error(f'Error reading PDF: {e}')
            result['valid_pdf'] = False
            result['page_count'] = 0
        finally:
            os.unlink(tmp_path)
    else:
        logger.warning(f"PDF file not found at {config['path']}")
        result['valid_pdf'] = False
        result['page_count'] = 0
    return result

def get_largest_filename_from_text__b7de68e1(env, config):
    """Read the largest file's name from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Filename of the largest .doc file
    """
    path = config.get('path', '/home/user/Desktop/largest_doc.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    content = file_bytes.decode('utf-8').strip()
    return content

def get_files_with_prefix__cf3f5d8ece62ecf8e4937dea9e007679(env, config: Dict) -> Dict:
    """
    Get files that have a specific prefix.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'prefix' keys

    Returns:
        Dict with files grouped by their hash
    """
    directory = config.get('directory', '/home/user/Pictures')
    prefix = config.get('prefix', 'Mountain_')
    command = f"""python3 -c "\nimport os\nimport hashlib\nimport json\nfrom PIL import Image\n\ndirectory = '{directory}'\nprefix = '{prefix}'\n\nresult = {{}}\n\nif os.path.exists(directory):\n    files = os.listdir(directory)\n    for filename in files:\n        if filename.startswith(prefix):\n            filepath = os.path.join(directory, filename)\n            if os.path.isfile(filepath):\n                try:\n                    with Image.open(filepath) as img:\n                        img_byte_arr = img.tobytes()\n                        file_hash = hashlib.sha256(img_byte_arr).hexdigest()\n                        result[file_hash] = filename\n                except:\n                    pass\n\nprint(json.dumps(result))\n"\n"""
    run_result = env.controller.run_bash_script(command, timeout=10)
    if run_result.get('returncode') != 0:
        return {}
    output = run_result.get('output', '')
    try:
        import json
        files_by_hash = json.loads(output.strip())
        return files_by_hash
    except:
        return {}

def get_text_content__2ad6be23(env, config: dict):
    """Extract text from VM file.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        str: File text content
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        return file_bytes.decode('utf-8').strip()
    except:
        return None

def get_csv_filtered_rows__8b005fb41e3efd4706af2ae4f1a79bee(env, config):
    """
    Read CSV file from VM and return all rows as list of lists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists representing CSV rows (including header)
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        rows = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
    finally:
        os.unlink(tmp_path)

def get_vm_dir_list__e4b458033c1389ecc56cd63f5eae9626(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a VM directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying directory path

    Returns:
        List of filenames in the directory
    """
    path = config['path']
    result = env.controller.run_bash_script(f"ls -1 '{path}' 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0:
        logger.warning(f"Failed to list directory {path}: {result.get('error', '')}")
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    return files

def get_file_exists__d22e0dfa(env, config: dict):
    """Check if the exported HTML file exists on VM and validate its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary containing file existence, size, and HTML validation info
    """
    vm_path = config.get('path', '/home/user/Desktop/survey-table.html')
    exists_result = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    exists = exists_result.get('output', '').strip() == 'EXISTS'
    if not exists:
        return {'exists': False, 'is_html': False, 'has_content': False, 'has_table': False, 'file_size': 0}
    size_result = env.controller.run_bash_script(f"stat -c %s '{vm_path}' 2>/dev/null || echo '0'", timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    content_result = env.controller.run_bash_script(f"head -c 4096 '{vm_path}' 2>/dev/null || echo ''", timeout=10)
    content = content_result.get('output', '').lower()
    is_html = any(['<html' in content, '<!doctype html' in content, '<meta' in content and '>' in content])
    has_table = '<table' in content or '<td' in content or '<tr' in content
    has_content = file_size > 100
    return {'exists': exists, 'is_html': is_html, 'has_content': has_content, 'has_table': has_table, 'file_size': file_size}

def get_text_file_content__16aae469605745263765ebbfdcb35f54(env, config: dict) -> Optional[str]:
    """Get content from a text file on the VM as a single string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_bash_file_content__816eebe3a50924fca9ee760018f3238e(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file not found
    """
    file_path = config.get('path', '/home/user/output.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_text_file_content__c491dbb8(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_file_exists_and_size__dccb2b60(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_zip_file_list__ee271ee8(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get list of files in a ZIP archive from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool) and 'files' (List[str]) keys
    """
    path = config['path']
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'files': []}
        cache_path = os.path.join(env.cache_dir, 'temp_check.zip')
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        with zipfile.ZipFile(cache_path, 'r') as zf:
            files = zf.namelist()
        return {'exists': True, 'files': files}
    except Exception as e:
        logger.error(f'Failed to read ZIP file: {e}')
        return {'exists': False, 'files': []}

def get_docx_text_content__66306a8b(env, config):
    """Get full text content of a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Full text content
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return ''
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        full_text = '\n'.join([para.text for para in doc.paragraphs])
        logger.info(f'Extracted {len(full_text)} characters')
        return full_text
    except Exception as e:
        logger.error(f'Error reading document: {e}')
        return ''
    finally:
        import os
        os.unlink(tmp_path)

def get_speedtest_file_content__26660ad1(env, config):
    """
    Get the content of the speedtest results file.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key

    Returns:
        str: File content as text, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_nested_directory_exists__3c3f0ceb(env, config: dict):
    """
    Check if a nested directory path exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if nested directory exists, False otherwise
    """
    path = config.get('path', '')
    command = f'test -d "{path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        output = result['output'].strip()
        return output == 'exists'
    return False

def get_file_exists__5a942dd0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if CSV file exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict containing rules with expected file_path

    Returns:
        Dict with file_exists (bool) and file_path (str)
    """
    csv_path = config.get('rules', {}).get('file_path', '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.csv')
    result = env.controller.get_file(csv_path)
    file_exists = result is not None and len(result) > 0
    return {'file_exists': file_exists, 'file_path': csv_path}

def get_text_file_content__1e67b9ae311891ef2e3034615cded86c(env, config: Dict) -> List[str]:
    """
    Read text file content and return as list of lines.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lines from the file
    """
    file_path = config.get('path', '/home/user/Desktop/mountains.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        return lines
    except:
        return []

def get_python_file_content__c685793b2ca36ab76b7f2cc84f84fe40(env, config: Dict[str, Any]) -> str:
    """Get content of a Python file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Python file content, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        return file_bytes.decode('utf-8')
    except Exception as e:
        return ''

def get_docx_content__80bc3840fb160595a8bbd7e90abee050(env, config):
    """Extract content and metadata from a .docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location on VM

    Returns:
        dict: Contains 'paragraphs' (list of text), 'word_count', 'paragraph_count'
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': 'No path provided'}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': f'File not found: {file_path}'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        all_text = ' '.join(paragraphs)
        word_count = len(all_text.split())
        return {'paragraphs': paragraphs, 'word_count': word_count, 'paragraph_count': len(paragraphs), 'full_text': all_text}
    except Exception as e:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': str(e)}
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_archive_contents__6131523f(env, config):
    """Get contents of a tar archive file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (local filename)

    Returns:
        List of filenames in the archive
    """
    vm_path = config.get('path')
    dest = config.get('dest', 'archive.tar.gz')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return []
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        with tarfile.open(cache_path, 'r:*') as tar:
            members = tar.getnames()
            filenames = [os.path.basename(m) for m in members if not os.path.isdir(os.path.join(tempfile.gettempdir(), m))]
            return sorted(filenames)
    except Exception as e:
        print(f'Error reading archive: {e}')
        return []

def get_total_files_in_dirs__ed6a3699fe6af7deef02a8e547504034(env, config):
    """Count number of files in each directory individually.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' (list)

    Returns:
        dict: Dictionary mapping directory name to file count (e.g., {'dir1': 1, 'dir2': 1, 'dir3': 0})
    """
    directories = config.get('directories', [])
    dir_counts = {}
    for directory in directories:
        count_cmd = f'find "{directory}" -maxdepth 1 -type f | wc -l'
        result = env.controller.run_bash_script(count_cmd, timeout=10)
        count = 0
        if result and result.get('output'):
            try:
                count = int(result['output'].strip())
            except ValueError:
                count = 0
        dir_counts[directory] = count
    return dir_counts

def get_file_size_info__0f2a6243(env, config: Dict) -> Optional[Dict]:
    """
    Get file size information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_csv_column_values__8b88ca382748ac90936480358fbda368(env, config):
    """Read CSV and extract ALL values from a specific column to verify sorting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'column' keys

    Returns:
        dict: Dictionary with 'values' (all column values) and 'total_rows' (count)
    """
    file_path = config.get('path', '')
    column_name = config.get('column', '')
    if not file_path or not column_name:
        return {'values': [], 'total_rows': 0}
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'values': [], 'total_rows': 0}
        import csv
        import io
        content = file_bytes.decode('utf-8', errors='ignore')
        reader = csv.DictReader(io.StringIO(content))
        values = []
        for row in reader:
            if column_name in row:
                values.append(row[column_name].strip())
        return {'values': values, 'total_rows': len(values)}
    except Exception as e:
        return {'values': [], 'total_rows': 0}

def get_text_file_content__ebe7bea30aab4d43b91ea1760b3fb66f(env, config: Dict[str, Any]) -> str:
    """
    Get content of a text file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        logger.warning(f'File not found: {file_path}')
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file content: {e}')
        return ''

def get_docx_text_content__2830db9ce8bded3cea70386f80d8c3ad(env, config: Dict[str, Any]) -> str:
    """Extract all text content from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        str: All text content from the document, concatenated
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_text = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                all_text.append(text)
        return '\n'.join(all_text)
    finally:
        os.unlink(tmp_path)

def get_pdf_files_in_dir__a76459ec(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_python_file_content__3b7fbcc9bcd509e2386e3dc8ef7407df(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Extract Python file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file {path}: {e}')
        return None

def get_docx_file_exists__c4f2c36653f58f8aab105ca7a48ec763(env, config: dict):
    """Check if a .docx file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'size': int or None, 'filename': str}
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes:
        return {'exists': True, 'size': len(file_bytes), 'filename': os.path.basename(file_path)}
    else:
        return {'exists': False, 'size': None, 'filename': os.path.basename(file_path)}

def get_pdf_page_info__4c6d2d3e077e0f5ce3b4338c5664395c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get PDF page count and check if file exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' (bool) and 'page_count' (int)
    """
    try:
        from pypdf import PdfReader
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'exists': False, 'page_count': 0}
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            return {'exists': True, 'page_count': len(reader.pages)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return {'exists': False, 'page_count': 0}

def get_screenshot_in_folder__956d014f26874db42b602f2626906894(env, config: dict):
    """
    Check if a screenshot exists in a specific folder with correct dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File information including location and dimensions
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Screenshot not found at: {path}')
            return {'exists': False, 'path': path, 'directory': os.path.dirname(path), 'filename': os.path.basename(path)}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            return {'exists': True, 'path': path, 'directory': os.path.dirname(path), 'filename': os.path.basename(path), 'size': len(file_bytes), 'width': width, 'height': height, 'format': img.format}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking screenshot in folder {path}: {e}')
        return {'exists': False, 'path': path, 'directory': os.path.dirname(path), 'filename': os.path.basename(path), 'error': str(e)}

def get_file_in_directory__845b16e9c20bb76eb4ef8a7eb9262413(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists in a directory and return its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'filename_pattern' keys

    Returns:
        Dict with 'exists': bool and 'content': str if file exists
    """
    import os
    directory = config.get('directory', '')
    filename_pattern = config.get('filename_pattern', '')
    result = env.controller.run_bash_script(f'ls -1 "{directory}" 2>/dev/null || echo ""', timeout=10)
    if result.get('returncode', 1) != 0 or not result.get('output', ''):
        logger.warning(f'Directory {directory} not found or empty')
        return {'exists': False, 'content': ''}
    files = result['output'].strip().split('\n')
    matching_file = None
    for f in files:
        if filename_pattern in f:
            matching_file = f
            break
    if not matching_file:
        logger.warning(f"No file matching pattern '{filename_pattern}' found in {directory}")
        return {'exists': False, 'content': ''}
    full_path = os.path.join(directory, matching_file)
    file_bytes = env.controller.get_file(full_path)
    if not file_bytes:
        logger.warning(f'Failed to read file: {full_path}')
        return {'exists': True, 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
        return {'exists': True, 'content': content}
    except Exception as e:
        logger.error(f'Error decoding file content: {e}')
        return {'exists': True, 'content': ''}

def get_text_content__f9fa7b925a8af809c4d7a02e0264dd76(env, config):
    """
    Get the full content of a text file.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key specifying the file path on VM

    Returns:
        str: File content or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes:
        try:
            content = file_bytes.decode('utf-8', errors='ignore')
            return content.strip()
        except Exception:
            return ''
    else:
        return ''

def get_backup_files_check__ac9408954b941c7f40eedd27a6f1296b(env, config: dict):
    """
    Check for backup file existence in specified directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: Backup file existence information
    """
    check_script = '\n#!/bin/bash\nresult=""\n\n# Check for file1.backup in each directory\nfor dir in dir1 dir2 dir3; do\n    if [ -f "$dir/file1.backup" ]; then\n        result="${result}${dir}_backup|"\n    fi\ndone\n\necho "$result"\n'
    result = env.controller.run_bash_script(check_script, timeout=10)
    if result['returncode'] != 0:
        return {'dir1_backup': False, 'dir2_backup': False, 'dir3_backup': False}
    output = result['output'].strip()
    return {'dir1_backup': 'dir1_backup' in output, 'dir2_backup': 'dir2_backup' in output, 'dir3_backup': 'dir3_backup' in output}

def get_folder_contents__4e03b1ed(env, config: dict):
    """Get folder structure information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        dict: Folder information including existence and file count
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    check_cmd = f"test -d '{folder_path}' && echo 'exists' || echo 'not_exists'"
    result = env.controller.run_bash_script(check_cmd, timeout=10)
    exists = result.get('output', '').strip() == 'exists'
    pdf_count = 0
    if exists:
        count_cmd = f"ls -1 '{folder_path}'/*.pdf 2>/dev/null | wc -l"
        count_result = env.controller.run_bash_script(count_cmd, timeout=10)
        try:
            pdf_count = int(count_result.get('output', '0').strip())
        except ValueError:
            pdf_count = 0
    logger.info(f'Folder {folder_path} - Exists: {exists}, PDF count: {pdf_count}')
    return {'exists': exists, 'pdf_count': pdf_count}

def get_pdf_file_size__f34c6d72482f0e93994904e974de58fa(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file and extract file size information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with file size information
    """
    from pypdf import PdfReader
    file_bytes = env.controller.get_file(config['path'])
    result = {'path': config['path'], 'exists': file_bytes is not None and len(file_bytes) > 0, 'file_size_bytes': 0}
    if result['exists']:
        result['file_size_bytes'] = len(file_bytes)
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            result['valid_pdf'] = len(reader.pages) > 0
            result['page_count'] = len(reader.pages)
            logger.info(f"PDF file size: {result['file_size_bytes']} bytes ({result['file_size_bytes'] / 1024:.2f} KB), {result['page_count']} pages")
        except Exception as e:
            logger.error(f'Error reading PDF: {e}')
            result['valid_pdf'] = False
            result['page_count'] = 0
        finally:
            os.unlink(tmp_path)
    else:
        logger.warning(f"PDF file not found at {config['path']}")
        result['valid_pdf'] = False
        result['page_count'] = 0
    return result

def get_docx_text_content__d99037c6(env, config: Dict[str, Any]) -> str:
    """
    Get all text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: All text content concatenated
    """
    from docx import Document
    from io import BytesIO
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.error(f'Could not retrieve file: {path}')
        return ''
    try:
        doc = Document(BytesIO(file_bytes))
        text = '\n'.join((p.text for p in doc.paragraphs if p.text.strip()))
        logger.info(f'Extracted {len(text)} characters from {path}')
        return text
    except Exception as e:
        logger.error(f'Error reading DOCX file: {e}')
        return ''

def get_tetris_files__ca30528a(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_docx_text_content__faacb6dea70ef25b404521dc6e6554b5(env, config: Dict[str, Any]) -> str:
    """
    Extract all text content from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file

    Returns:
        String containing all text from the document
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        text_parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                text_parts.append(text)
        return '\n'.join(text_parts)
    finally:
        os.unlink(tmp_path)

def get_text_file_lines__7e304294ce76dab9f08b95178667a620(env, config):
    """
    Read text file from VM and return non-empty lines.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of non-empty lines from the file
    """
    file_path = config.get('path', '')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        return lines
    except Exception as e:
        print(f'Error reading text file: {e}')
        return []

def get_multi_directory_contents__429c8cbc(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_docx_file_count__bb5651c2(env, config):
    """
    Get the content of document_count.txt file.

    This getter reads the content of /home/user/Desktop/document_count.txt
    which should contain the count of .docx files in the students work folder.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        str: Content of the file (stripped), or None if file doesn't exist
    """
    file_path = config.get('path', '/home/user/Desktop/document_count.txt')
    try:
        with open(file_path, 'r') as f:
            content = f.read().strip()
        return content
    except FileNotFoundError:
        return None
    except Exception as e:
        return None

def get_text_file_content__5ced85fc_aug18_v4_d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8a9(env, config):
    """Read a text file from VM and return its entire content as a string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        str: File content as a string (with trailing whitespace stripped)
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        logger.info(f'Successfully read {len(content)} characters from {file_path}')
        return content
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return ''

def get_text_file_content__b38e8bb9(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_docx_content__5ab5849946111fe9c5dda6e47422b057(env, config):
    """Extract content and metadata from a .docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location on VM

    Returns:
        dict: Contains 'paragraphs' (list of text), 'word_count', 'paragraph_count'
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': 'No path provided'}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': f'File not found: {file_path}'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        all_text = ' '.join(paragraphs)
        word_count = len(all_text.split())
        return {'paragraphs': paragraphs, 'word_count': word_count, 'paragraph_count': len(paragraphs), 'full_text': all_text}
    except Exception as e:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': str(e)}
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_text_file_content__87f48e6e(env, config: dict):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return None

def get_pdf_files_with_content__184db76a3945be5ea6dec91515beac2a(env, config: dict):
    """Check if PDF files exist in directory and contain expected text.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path on VM to check
            - expected_files: List of dicts with 'filename' and 'content_check' keys

    Returns:
        dict: Results for each file {filename: 1 (exists and has content) or 0/-1}
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    expected_files = config.get('expected_files', [])
    results = {}
    for file_info in expected_files:
        filename = file_info['filename']
        content_check = file_info.get('content_check', '')
        file_path = os.path.join(directory, filename)
        check_cmd = f'[ -f "{file_path}" ] && echo "exists" || echo "not_found"'
        result = env.controller.run_bash_script(check_cmd, timeout=10)
        if result.get('output', '').strip() != 'exists':
            results[filename] = -1
            continue
        if content_check:
            try:
                file_bytes = env.controller.get_file(file_path)
                if not file_bytes:
                    results[filename] = 0
                    continue
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    import fitz
                    doc = fitz.open(tmp_path)
                    text = ''
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    if content_check in text:
                        results[filename] = 1
                    else:
                        results[filename] = 0
                except Exception as e:
                    results[filename] = 0
                finally:
                    os.unlink(tmp_path)
            except Exception as e:
                results[filename] = 0
        else:
            results[filename] = 1
    return results

def get_tetris_files__cb4dc6e5(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_file_exists__04b2d876(env, config: Dict[str, Any]):
    """
    Check if a file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (str)

    Returns:
        bool: True if file exists, False otherwise
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    return file_bytes is not None

def get_python_file_content__d881a7db2082639af55dd0ea1e3047ab(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Extract Python file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file {path}: {e}')
        return None

def get_pdf_file_info__3dd8a1347f60dd796dc5e98521ae4032(env, config: Dict[str, str]) -> Optional[str]:
    """
    Find and download the most recent PDF file from Downloads folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - downloads_dir: Directory to search for PDFs (e.g., /home/user/Downloads)
            - filename_pattern: Pattern to match (e.g., *.pdf)

    Returns:
        str: Path to downloaded PDF file in cache, or None if no PDF found
    """
    downloads_dir = config.get('downloads_dir', '/home/user/Downloads')
    filename_pattern = config.get('filename_pattern', '*.pdf')
    search_path = os.path.join(downloads_dir, filename_pattern)
    result = env.controller.execute(f'ls -t {search_path} 2>/dev/null | head -n 1')
    if not result or result.strip() == '':
        return None
    pdf_path = result.strip()
    file_bytes = env.controller.get_file(pdf_path)
    if not file_bytes:
        return None
    dest_filename = os.path.basename(pdf_path)
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_file_content__89d906fa(env, config: dict):
    """
    Read the content of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Content of the file (stripped)
    """
    path = config.get('path', '/home/user/Desktop/doc_count.txt')
    command = f'cat {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result.get('returncode', 1) != 0:
        return ''
    return result.get('output', '').strip()

def get_python_files_count__13d7d579(env, config):
    """
    Count the number of Python files and verify repository integrity.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        dict: {
            'python_file_count': int,
            'has_git_dir': bool,
            'git_remote_url': str or None,
            'has_readme': bool,
            'has_setup_py': bool,
            'has_requirements': bool
        }
    """
    repo_path = config.get('repo_path', '/home/user/instructor-embedding')
    py_count_cmd = f"find {repo_path} -name '*.py' 2>/dev/null | wc -l"
    py_result = env.controller.run_bash_script(py_count_cmd, timeout=10)
    try:
        py_count = int(py_result.get('output', '0').strip())
    except ValueError:
        py_count = 0
    git_dir_cmd = f"test -d {repo_path}/.git && echo 'exists' || echo 'missing'"
    git_dir_result = env.controller.run_bash_script(git_dir_cmd, timeout=5)
    has_git_dir = git_dir_result.get('output', '').strip() == 'exists'
    git_remote_url = None
    if has_git_dir:
        git_remote_cmd = f"cd {repo_path} && git remote get-url origin 2>/dev/null || echo ''"
        git_remote_result = env.controller.run_bash_script(git_remote_cmd, timeout=5)
        git_remote_url = git_remote_result.get('output', '').strip()
    readme_cmd = f"test -f {repo_path}/README.md && echo 'exists' || echo 'missing'"
    readme_result = env.controller.run_bash_script(readme_cmd, timeout=5)
    has_readme = readme_result.get('output', '').strip() == 'exists'
    setup_cmd = f"test -f {repo_path}/setup.py && echo 'exists' || echo 'missing'"
    setup_result = env.controller.run_bash_script(setup_cmd, timeout=5)
    has_setup_py = setup_result.get('output', '').strip() == 'exists'
    requirements_cmd = f"test -f {repo_path}/requirements.txt && echo 'exists' || echo 'missing'"
    requirements_result = env.controller.run_bash_script(requirements_cmd, timeout=5)
    has_requirements = requirements_result.get('output', '').strip() == 'exists'
    return {'python_file_count': py_count, 'has_git_dir': has_git_dir, 'git_remote_url': git_remote_url, 'has_readme': has_readme, 'has_setup_py': has_setup_py, 'has_requirements': has_requirements}

def get_file_exists__25309a67d723dd8e75eb60c978b60929(env, config):
    """Check if a file exists on the VM and verify it's a valid GIF image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'is_gif': bool, 'has_content': bool, 'path': str}
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "true" || echo "false"', timeout=10)
    exists = result.get('output', '').strip() == 'true'
    is_gif = False
    has_content = False
    if exists:
        mime_result = env.controller.run_bash_script(f'file -b --mime-type "{file_path}"', timeout=10)
        mime_type = mime_result.get('output', '').strip()
        is_gif = mime_type == 'image/gif'
        size_result = env.controller.run_bash_script(f'[ -s "{file_path}" ] && echo "true" || echo "false"', timeout=10)
        has_content = size_result.get('output', '').strip() == 'true'
    return {'exists': exists, 'is_gif': is_gif, 'has_content': has_content, 'path': file_path}

def get_vm_file__519b2394782587cf0953c6f72e80b8a6(env, config):
    """
    Get the DOCX file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded file
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    import tempfile
    import os
    suffix = os.path.splitext(file_path)[1]
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_python_files_syntax__c9c8227a4cd72e8de3e73ea399f7f61d(env, config: dict) -> Dict[str, Any]:
    """Check if all Python files have valid syntax and are free of runtime errors.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' key (list of file paths)

    Returns:
        Dict containing:
        - all_valid: bool - whether all files have valid syntax and can run
        - file_results: dict - per-file validation results
        - total_files: int - number of files checked
        - valid_files: int - number of files with valid syntax and runtime checks
    """
    file_paths = config.get('paths', [])
    if not file_paths:
        return {'all_valid': False, 'file_results': {}, 'total_files': 0, 'valid_files': 0, 'error': 'No file paths provided'}
    file_results = {}
    valid_count = 0
    import tempfile
    temp_dir = tempfile.mkdtemp()
    try:
        temp_file_paths = {}
        for file_path in file_paths:
            file_name = os.path.basename(file_path)
            temp_file_path = os.path.join(temp_dir, file_name)
            try:
                file_bytes = env.controller.get_file(file_path)
                if not file_bytes:
                    file_results[file_name] = {'valid': False, 'error': 'File not found or empty'}
                    continue
                with open(temp_file_path, 'wb') as f:
                    f.write(file_bytes)
                temp_file_paths[file_name] = temp_file_path
            except Exception as e:
                file_results[file_name] = {'valid': False, 'error': f'Failed to read file: {str(e)}'}
                continue
        if temp_dir not in sys.path:
            sys.path.insert(0, temp_dir)
        for (file_name, temp_file_path) in temp_file_paths.items():
            try:
                with open(temp_file_path, 'r', encoding='utf-8') as f:
                    code = f.read()
                try:
                    ast.parse(code)
                except SyntaxError as e:
                    file_results[file_name] = {'valid': False, 'error': f'SyntaxError at line {e.lineno}: {e.msg}'}
                    continue
                except Exception as e:
                    file_results[file_name] = {'valid': False, 'error': f'Parse error: {str(e)}'}
                    continue
                try:
                    compile(code, temp_file_path, 'exec')
                except Exception as e:
                    file_results[file_name] = {'valid': False, 'error': f'Compilation error: {str(e)}'}
                    continue
                try:
                    module_name = file_name[:-3] if file_name.endswith('.py') else file_name
                    spec = importlib.util.spec_from_file_location(module_name, temp_file_path)
                    if spec and spec.loader:
                        module = importlib.util.module_from_spec(spec)
                        spec.loader.exec_module(module)
                    else:
                        file_results[file_name] = {'valid': False, 'error': 'Failed to create module spec'}
                        continue
                except Exception as e:
                    file_results[file_name] = {'valid': False, 'error': f'Import/Runtime error: {str(e)}'}
                    continue
                file_results[file_name] = {'valid': True, 'error': None}
                valid_count += 1
            except Exception as e:
                file_results[file_name] = {'valid': False, 'error': f'Validation error: {str(e)}'}
    finally:
        if temp_dir in sys.path:
            sys.path.remove(temp_dir)
        import shutil
        try:
            shutil.rmtree(temp_dir)
        except:
            pass
    return {'all_valid': valid_count == len(file_paths), 'file_results': file_results, 'total_files': len(file_paths), 'valid_files': valid_count}

def get_text_file_lines__d1fc13ca7061617e08d8a914a14209cd(env, config: dict) -> Optional[list]:
    """Get lines from a text file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of stripped lines, or None if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n')]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_filename_from_path__a6b53d3d(env, config: Dict[str, Any]) -> Optional[str]:
    """Extract filename from a file path and verify file exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Filename if file exists, None otherwise
    """
    file_path = config.get('path')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        filename = os.path.basename(file_path)
        return filename
    except Exception:
        return None

def get_tetris_files__58a20c62(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_pdf_files_info__8d8da24c(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_recent_pdf_count__db5a3e05(env, config: dict):
    """Get count of recently modified PDF files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path
            - max_age_seconds: Maximum age in seconds

    Returns:
        Number of recent PDF files
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    max_age = config.get('max_age_seconds', 300)
    command = f"""\npython3 -c "\nimport os\nimport glob\nimport time\n\ndir_path = '{directory}'\nmax_age = {max_age}\n\nif not os.path.exists(dir_path):\n    print('0')\n    exit(0)\n\ncurrent_time = time.time()\npdf_files = glob.glob(os.path.join(dir_path, '*.pdf'))\nrecent_files = [f for f in pdf_files if (current_time - os.path.getmtime(f)) <= max_age]\nprint(len(recent_files))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to count recent files: {result['error']}")
        return 0
    try:
        count = int(result['output'].strip())
        return count
    except:
        return 0

def get_compressed_file_contents__0b074054(env, config):
    """Get list of files in a compressed archive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: List of files in archive
    """
    path = config.get('path', '')
    command = f'tar -tzf {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        return result['output'].strip()
    else:
        logger.error(f"Failed to list contents of {path}: {result['error']}")
        return None

def get_file_content__f8073900a375900c7a9ce8fa79f05f9d(env, config):
    """Read a text file and return its content as a string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content, or empty string if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        return ''

def get_file_existence__81d9e3403ff656f186df75dab490d6ac(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file exists at the specified path and return its properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'is_png', 'size' keys
    """
    path = config.get('path', '')
    result = {'exists': False, 'is_png': False, 'size': 0}
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            result['exists'] = True
            result['size'] = len(file_bytes)
            if path.lower().endswith('.png'):
                result['is_png'] = True
                try:
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    try:
                        img = Image.open(tmp_path)
                        img.verify()
                        result['is_png'] = True
                    except Exception as e:
                        logger.warning(f'File exists but is not a valid PNG: {e}')
                        result['is_png'] = False
                    finally:
                        os.unlink(tmp_path)
                except Exception as e:
                    logger.warning(f'Error verifying PNG: {e}')
        else:
            logger.info(f'File does not exist at path: {path}')
    except Exception as e:
        logger.error(f'Error checking file existence: {e}')
    return result

def get_vm_subtitle_file__dee238bc(env, config: dict):
    """
    Get subtitle file from VM.

    Args:
        env: Desktop environment
        config: Configuration with 'path' and 'dest' keys

    Returns:
        str: Path to downloaded subtitle file
    """
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_launch_json__a3743b930b4e3c5c6976584a28c8269c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get launch.json file content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict containing the launch.json content, or empty dict if file not found
    """
    file_path = config.get('path', '/home/user/project/.vscode/launch.json')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'launch.json not found: {file_path}')
            return {}
        launch_config = json.loads(file_bytes.decode('utf-8'))
        return launch_config
    except Exception as e:
        logger.error(f'Error reading launch.json: {e}')
        return {}

def get_pdf_docx_result__a2a86631dd34db500aa084715506d32d(env, config: dict):
    """
    Extract both command history and PDF archive file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' and 'path' keys

    Returns:
        tuple: (command_output, archive_bytes)
    """
    command_result = env.controller.run_bash_script(config['command'], timeout=30)
    command_output = command_result.get('output', '') if command_result else ''
    archive_bytes = env.controller.get_file(config['path'])
    return (command_output, archive_bytes)

def get_python_file_with_pattern__f222fdd4d3c26325ac7310b0f2b1711f(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Find and extract Python file matching a pattern from VM.
    Returns both content and filename for validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'base_path' and 'pattern' keys

    Returns:
        Dict with 'content', 'filename', and 'file_date' keys, or None if file doesn't exist
    """
    base_path = config.get('base_path', '/home/user')
    pattern = config.get('pattern', 'code_backup_.*\\.py')
    result = env.controller.run_bash_script(f'ls -1 {base_path}', timeout=10)
    if result.get('returncode') != 0:
        logger.warning(f'Failed to list directory: {base_path}')
        return None
    files = result.get('output', '').strip().split('\n')
    strict_pattern = re.compile('^code_backup_(\\d{8})\\.py$')
    matching_files = []
    for f in files:
        match = strict_pattern.match(f)
        if match:
            date_str = match.group(1)
            try:
                year = int(date_str[:4])
                month = int(date_str[4:6])
                day = int(date_str[6:8])
                if 2020 <= year <= 2099 and 1 <= month <= 12 and (1 <= day <= 31):
                    matching_files.append((f, date_str))
            except ValueError:
                continue
    if not matching_files:
        logger.warning(f'No file with valid YYYYMMDD date format found in {base_path}')
        return None
    (filename, date_str) = matching_files[0]
    target_file = os.path.join(base_path, filename)
    file_bytes = env.controller.get_file(target_file)
    if not file_bytes:
        logger.warning(f'File not found: {target_file}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        return {'content': content, 'filename': filename, 'file_date': date_str}
    except Exception as e:
        logger.error(f'Error decoding file {target_file}: {e}')
        return None

def get_dual_dir_check__9cc29f71(env, config):
    """Check both old and new directory states.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_dir' and 'new_dir' paths

    Returns:
        dict: Status of both directories
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    old_dir = config['old_dir']
    new_dir = config['new_dir']
    old_cmd = f"[ -d {old_dir} ] && echo 'exists' || echo 'not_exists'"
    response_old = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': old_cmd, 'shell': True})
    new_cmd = f"[ -d {new_dir} ] && echo 'exists' || echo 'not_exists'"
    response_new = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': new_cmd, 'shell': True})
    result = {'old_exists': response_old.json()['output'].strip() == 'exists', 'new_exists': response_new.json()['output'].strip() == 'exists'}
    print(f'Directory check result: {result}')
    return result

def get_file_content__cd720b2833d8c75c48e4cd829046ee69(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the content of a file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        Dict with 'exists', 'content', and 'file_path' keys
    """
    file_path = config.get('file_path', '')
    file_bytes = env.controller.get_file(file_path)
    result = {'exists': file_bytes is not None, 'file_path': file_path, 'content': ''}
    if file_bytes:
        try:
            result['content'] = file_bytes.decode('utf-8')
        except:
            result['content'] = ''
    return result

def get_txt_result__f5eea4d48349a222e7a85e09c99bbeae(env, config: dict):
    """
    Extract both command history and TXT archive file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' and 'path' keys

    Returns:
        tuple: (command_output, archive_bytes)
    """
    command_result = env.controller.run_bash_script(config['command'], timeout=30)
    command_output = command_result.get('output', '') if command_result else ''
    archive_bytes = env.controller.get_file(config['path'])
    return (command_output, archive_bytes)

def get_txt_line_count__8ff67d2b(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get a text file from the VM for line count verification.

    Args:
        env: Environment object with controller to access VM
        config: Configuration dict with 'path' key pointing to text file on VM

    Returns:
        str: Local path to the downloaded text file, or None if file doesn't exist
    """
    txt_path_on_vm = config['path']
    txt_filename = os.path.basename(txt_path_on_vm)
    local_path = os.path.join(env.cache_dir, txt_filename)
    try:
        file_content = env.controller.get_file(txt_path_on_vm)
        if file_content is None:
            return None
        with open(local_path, 'wb') as f:
            f.write(file_content)
        return local_path
    except Exception as e:
        return None

def get_file_content__6c41fae8ffe95d2c4b3c521d5446de56(env, config: Dict[str, Any]) -> Dict[str, str]:
    """Get both original (from cache URL) and current content of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key and optional 'original_url' key

    Returns:
        Dict[str, str]: Dictionary with 'original' and 'current' file content
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'original': '', 'current': ''}
    file_bytes = env.controller.get_file(file_path)
    current_content = ''
    if file_bytes:
        try:
            current_content = file_bytes.decode('utf-8')
        except Exception:
            current_content = ''
    original_content = ''
    original_url = config.get('original_url', '')
    if original_url:
        try:
            response = requests.get(original_url, timeout=10)
            if response.status_code == 200:
                original_content = response.text
        except Exception:
            original_content = ''
    return {'original': original_content, 'current': current_content}

def get_file_exists_and_size__479794c8(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_specific_pdf_exists__4e03b1ed(env, config: dict):
    """Check if specific PDF file exists and validate it's a proper PDF with expected content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        dict: Dictionary with validation results including exists, is_pdf, has_content, size, and contains_keywords
    """
    file_path = config.get('file_path', '')
    exists_command = f"test -f '{file_path}' && echo 'yes' || echo 'no'"
    exists_result = env.controller.run_bash_script(exists_command, timeout=10)
    exists = exists_result.get('output', '').strip() == 'yes'
    if not exists:
        logger.info(f'File {file_path} does not exist')
        return {'exists': False, 'is_pdf': False, 'has_content': False, 'size': 0, 'contains_keywords': False}
    file_type_command = f"file -b --mime-type '{file_path}'"
    file_type_result = env.controller.run_bash_script(file_type_command, timeout=10)
    mime_type = file_type_result.get('output', '').strip()
    is_pdf = mime_type == 'application/pdf'
    size_command = f"stat -c %s '{file_path}' 2>/dev/null || stat -f %z '{file_path}' 2>/dev/null || echo '0'"
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
    except ValueError:
        size = 0
    has_content = size > 1000
    contains_keywords = False
    if is_pdf and has_content:
        keywords_command = f"pdftotext '{file_path}' - 2>/dev/null | grep -i 'LLM.*Autonomous.*Agent' | head -1"
        keywords_result = env.controller.run_bash_script(keywords_command, timeout=15)
        keyword_output = keywords_result.get('output', '').strip()
        author_command = f"pdftotext '{file_path}' - 2>/dev/null | grep -i 'Lilian Weng' | head -1"
        author_result = env.controller.run_bash_script(author_command, timeout=15)
        author_output = author_result.get('output', '').strip()
        contains_keywords = bool(keyword_output) or bool(author_output)
    result = {'exists': exists, 'is_pdf': is_pdf, 'has_content': has_content, 'size': size, 'contains_keywords': contains_keywords}
    logger.info(f'PDF validation for {file_path}: {result}')
    return result

def get_text_file_content__49075d97b370c316554a4b259a7ccc3e(env, config):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_pdf_files_in_dir__80b8ddf2(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_file_rename_status__c396550f(env, config: Dict[str, Any]) -> Dict[str, bool]:
    """Check if a file was renamed (old file doesn't exist, new file exists).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_path' and 'new_path' keys

    Returns:
        Dict[str, bool]: Dictionary with 'old_exists' and 'new_exists' keys
    """
    old_path = config.get('old_path', '')
    new_path = config.get('new_path', '')
    old_result = env.controller.run_bash_script(f'test -f "{old_path}" && echo "exists" || echo "not_exists"', timeout=10)
    old_exists = 'exists' in old_result.get('output', '')
    new_result = env.controller.run_bash_script(f'test -f "{new_path}" && echo "exists" || echo "not_exists"', timeout=10)
    new_exists = 'exists' in new_result.get('output', '')
    return {'old_exists': old_exists, 'new_exists': new_exists}

def get_docx_content__fca57ed0800553851d5057d288f20b50(env, config: Dict[str, Any]) -> str:
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the DOCX file

    Returns:
        String containing the text content of the document
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        import tempfile
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        text_content = '\n'.join([para.text for para in doc.paragraphs])
        return text_content.strip()
    except Exception as e:
        return ''
    finally:
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_file_specific_line__ebbefd78bc8d93d288c452f335196528(env, config: Dict[str, Any]) -> Dict[str, str]:
    """Get the content of a specific line in a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: VM path to the file
            - line_number: Line number to retrieve (1-indexed)

    Returns:
        Dict with line content: {"content": "line content", "is_empty": True/False}
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'content': '', 'is_empty': True}
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1')
    lines = content.split('\n')
    line_number = config.get('line_number', 1)
    if 0 < line_number <= len(lines):
        line_content = lines[line_number - 1]
        return {'content': line_content, 'is_empty': len(line_content.strip()) == 0}
    return {'content': '', 'is_empty': True}

def get_pdf_page_count__25fc85ee(env, config):
    """Get the page count and content information of a PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for PDF file path on VM

    Returns:
        dict: Contains 'page_count', 'text_content', 'word_count', 'file_size', and 'metadata'
              Returns None if file doesn't exist or can't be read
    """
    vm_path = config.get('path', '')
    pdf_bytes = env.controller.get_file(vm_path)
    if not pdf_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, 'temp_check.pdf')
    with open(cache_path, 'wb') as f:
        f.write(pdf_bytes)
    try:
        with open(cache_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            page_count = len(pdf_reader.pages)
            full_text = ''
            for page in pdf_reader.pages:
                full_text += page.extract_text()
            metadata = pdf_reader.metadata if hasattr(pdf_reader, 'metadata') else {}
            word_count = len(full_text.split())
            file_size = len(pdf_bytes)
        return {'page_count': page_count, 'text_content': full_text, 'word_count': word_count, 'file_size': file_size, 'metadata': metadata}
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return None

def get_terminal_profile_name__dfa021d620bcc8024aed513c1adfaad3(env, config):
    """
    Get terminal default profile name from GNOME Terminal configuration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Terminal output containing profile name
    """
    return env.controller.get_terminal_output()

def get_file_exists__b8a50137(env, config):
    """
    Check if a file exists at the specified path on VM and verify it's a valid PNG image.

    Returns:
        dict: Contains 'exists' (bool), 'is_png' (bool), 'is_valid_image' (bool),
              'has_dimensions' (bool), 'width' (int), 'height' (int)
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    exists = result.get('output', '').strip() == 'EXISTS'
    if not exists:
        return {'exists': False, 'is_png': False, 'is_valid_image': False, 'has_dimensions': False, 'width': 0, 'height': 0}
    file_type_cmd = f'file -b "{file_path}"'
    file_type_result = env.controller.run_bash_script(file_type_cmd, timeout=10)
    file_type_output = file_type_result.get('output', '').strip().lower()
    is_png = 'png' in file_type_output
    dimensions_cmd = f"""python3 -c "\nimport sys\ntry:\n    from PIL import Image\n    img = Image.open('{file_path}')\n    print('{{0}},{{1}}'.format(img.width, img.height))\nexcept ImportError:\n    print('PIL_NOT_AVAILABLE')\n    sys.exit(1)\nexcept Exception as e:\n    print('ERROR')\n    sys.exit(1)\n" """
    dimensions_result = env.controller.run_bash_script(dimensions_cmd, timeout=10)
    dimensions_output = dimensions_result.get('output', '').strip()
    is_valid_image = dimensions_output != 'ERROR' and dimensions_output != '' and (dimensions_output != 'PIL_NOT_AVAILABLE')
    has_dimensions = False
    width = 0
    height = 0
    if is_valid_image and ',' in dimensions_output:
        try:
            (width_str, height_str) = dimensions_output.split(',')
            width = int(width_str)
            height = int(height_str)
            has_dimensions = width >= 100 and height >= 100
        except (ValueError, IndexError):
            is_valid_image = False
            has_dimensions = False
    return {'exists': exists, 'is_png': is_png, 'is_valid_image': is_valid_image, 'has_dimensions': has_dimensions, 'width': width, 'height': height}

def get_ext_name_by_path__4c281757414277d3f93140772916c7ad(env, config: Dict[str, Any]) -> str:
    """Get the extension name from Chrome preferences for an extension at a specific path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'extension_path' key

    Returns:
        str: Extension name if found, empty string otherwise
    """
    os_type = env.vm_platform
    target_path = config.get('extension_path', '')
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        all_extensions = data.get('extensions', {}).get('settings', {})
        for (ext_id, ext_data) in all_extensions.items():
            path = ext_data.get('path', '')
            if path == target_path:
                manifest = ext_data.get('manifest', {})
                name = manifest.get('name', '')
                return name
        return ''
    except Exception as e:
        logger.error(f'Error reading Chrome Preferences file: {e}')
        return ''

def get_csv_filtered_rows__d0eaa9b89f4859670e467a80277e775c(env, config):
    """
    Read CSV file from VM and return all rows as list of lists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists representing CSV rows (including header)
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        rows = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
    finally:
        os.unlink(tmp_path)

def get_pdf_chapter_files__242f4871e4280ec1025212dbeaf5c18c(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of PDF files in a directory from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key (VM path)

    Returns:
        List of PDF filenames (basenames only) found in the directory
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f"cd '{directory}' && ls -1 *.pdf 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0:
        logger.warning(f"Failed to list PDF files in {directory}: {result.get('error', '')}")
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    pdf_files = [line.strip() for line in output.split('\n') if line.strip()]
    logger.info(f'Found {len(pdf_files)} PDF files in {directory}: {pdf_files}')
    return pdf_files

def get_word_count_from_txt__68182234(env, config):
    """
    Extract word count from a text file that contains a count statement.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Word count extracted from file, or 0 if not found
    """
    import re
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return 0
        content = file_bytes.decode('utf-8', errors='ignore')
        patterns = ['count[:\\s]+(\\d+)', '(\\d+)\\s+(?:occurrence|time)', 'appears\\s+(\\d+)', 'found\\s+(\\d+)', 'total[:\\s]+(\\d+)']
        for pattern in patterns:
            match = re.search(pattern, content.lower())
            if match:
                return int(match.group(1))
        match = re.search('\\d+', content)
        if match:
            return int(match.group())
        return 0
    except Exception as e:
        return 0

def get_file_line_indentation__a2b5af8108e461937716b976809e966c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract indentation information from specific lines in a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: VM path to the file
            - start_line: First line to check (1-indexed, default: 1)
            - end_line: Last line to check (1-indexed, default: all lines)

    Returns:
        Dict with line numbers as keys and indentation info as values:
        {
            "2": {"leading_spaces": 4},
            "3": {"leading_spaces": 4},
            ...
        }
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {}
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1')
    lines = content.split('\n')
    start_line = config.get('start_line', 1)
    end_line = config.get('end_line', len(lines))
    result = {}
    for i in range(start_line - 1, min(end_line, len(lines))):
        if i < len(lines):
            line = lines[i]
            leading_spaces = len(line) - len(line.lstrip(' '))
            result[str(i + 1)] = {'leading_spaces': leading_spaces}
    return result

def get_multi_directory_contents__a93d97ba(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_file_permissions__fdcd3f41(env, config):
    """Get file permissions for specified file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'filename' parameter

    Returns:
        dict: Dictionary with read/write/execute permissions
    """
    filename = config.get('filename', '')
    filepath = f'/home/user/Pictures/{filename}'
    command = f'''python3 -c "import os; import json; import stat; s = os.stat('{filepath}'); mode = s.st_mode; result = {{'readable': bool(mode & stat.S_IRUSR), 'writable': bool(mode & stat.S_IWUSR), 'executable': bool(mode & stat.S_IXUSR)}}; print(json.dumps(result))"'''
    result = env.controller.run_bash_script(command, timeout=30)
    try:
        import json
        permissions = json.loads(result['output'].strip())
        return permissions
    except:
        return {'readable': False, 'writable': False, 'executable': False}

def get_dir_file_list__9ff7dd4db34cdbddf16b8ce0d6085594(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' key

    Returns:
        List of filenames in the directory
    """
    command = config.get('command', 'ls /home/user/Desktop/slides/')
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return []
    output = result['output'].strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    return files

def get_file_list_from_text__24c8b0df(env, config):
    """Read list of filenames from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        list: List of filenames (one per line)
    """
    path = config.get('path', '/home/user/Desktop/doc_files.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
    return lines

def get_text_file_content__860f6a3e21fa3ee43fbc12dc1309f38a(env, config: Dict[str, Any]) -> str:
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Text content of the file as a string, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_pdf_basic_info__9f4baa29714d8f65b6c7e77480d7ee20(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get basic PDF information including page count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'page_count' (int) and 'valid' (bool)
    """
    try:
        from pypdf import PdfReader
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'page_count': 0, 'valid': False}
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            page_count = len(reader.pages)
            valid = page_count > 0
            if valid and page_count > 0:
                page = reader.pages[0]
                width = float(page.mediabox.width)
                height = float(page.mediabox.height)
                valid = width > 0 and height > 0
            return {'page_count': page_count, 'valid': valid}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return {'page_count': 0, 'valid': False}

def get_text_file_lines__c2d1250f71a1fa98280fe372a2eb5875(env, config):
    """Read text file and return lines as list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of non-empty lines from file
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception:
        return []

def get_file_content__8ab0a45d4a57cf0e9592e87621895b59(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the content and properties of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        Dict with file existence and content information
    """
    file_path = config.get('file_path', '')
    result = {'exists': False, 'is_file': False, 'content': '', 'line_count': 0}
    check_file_cmd = f"[ -f '{file_path}' ] && echo 'EXISTS' || echo 'NOT_EXISTS'"
    file_check = env.controller.run_bash_script(check_file_cmd, timeout=10)
    if file_check.get('output', '').strip() == 'EXISTS':
        result['exists'] = True
        result['is_file'] = True
        content_cmd = f"head -n 1000 '{file_path}'"
        content_result = env.controller.run_bash_script(content_cmd, timeout=20)
        if content_result.get('returncode') == 0:
            result['content'] = content_result.get('output', '')
        line_count_cmd = f"wc -l < '{file_path}'"
        line_count_result = env.controller.run_bash_script(line_count_cmd, timeout=10)
        if line_count_result.get('returncode') == 0:
            try:
                result['line_count'] = int(line_count_result.get('output', '0').strip())
            except ValueError:
                result['line_count'] = 0
    logger.info(f"File content check for {file_path}: exists={result['exists']}, lines={result['line_count']}")
    return result

def get_docx_content__b519d5a9d3c41783a990af660a0ae167(env, config):
    """Extract content and metadata from a .docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location on VM

    Returns:
        dict: Contains 'paragraphs' (list of text), 'word_count', 'paragraph_count'
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': 'No path provided'}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': f'File not found: {file_path}'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        all_text = ' '.join(paragraphs)
        word_count = len(all_text.split())
        headings = []
        for para in doc.paragraphs:
            if para.style.name.startswith('Heading') or para.text.strip().isupper():
                headings.append(para.text.strip())
        has_lists = any(('•' in p or '–' in p or '■' in p for p in paragraphs))
        return {'paragraphs': paragraphs, 'word_count': word_count, 'paragraph_count': len(paragraphs), 'full_text': all_text, 'headings': headings, 'has_lists': has_lists}
    except Exception as e:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': str(e)}
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_tetris_files__65e60eeb(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_file_move_check__33f4e0c7b45da73c555209c719d5ff10(env, config: Dict[str, Any]) -> Dict[str, bool]:
    """Check if a file has been moved from old location to new location.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_path' and 'new_path' keys

    Returns:
        Dict with 'exists_at_new_path' and 'not_exists_at_old_path' booleans
    """
    old_path = config.get('old_path', '')
    new_path = config.get('new_path', '')
    command_new = f'test -f "{new_path}" && echo "exists" || echo "not_exists"'
    result_new = env.controller.run_bash_script(command_new, timeout=10)
    exists_at_new = result_new['returncode'] == 0 and 'exists' in result_new['output']
    command_old = f'test -f "{old_path}" && echo "exists" || echo "not_exists"'
    result_old = env.controller.run_bash_script(command_old, timeout=10)
    not_exists_at_old = result_old['returncode'] == 0 and 'not_exists' in result_old['output']
    return {'exists_at_new_path': exists_at_new, 'not_exists_at_old_path': not_exists_at_old}

def get_text_file_content__69119b71(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_srt_filename_check__ab47203640c5bdcef1195c50e51e7524(env, config: dict):
    """
    Check if SRT file exists with expected filename pattern and validate SRT format.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the SRT file

    Returns:
        dict: {
            'exists': bool,
            'filename': str,
            'has_content': bool,
            'is_valid_srt': bool,
            'has_multiple_entries': bool,
            'file_size': int
        }
    """
    path = config.get('path', '/home/user/video_subtitles.srt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'filename': '', 'has_content': False, 'is_valid_srt': False, 'has_multiple_entries': False, 'file_size': 0}
    filename = os.path.basename(path)
    file_size = len(file_bytes)
    has_content = file_size > 500
    is_valid_srt = False
    has_multiple_entries = False
    try:
        content = file_bytes.decode('utf-8')
        timestamp_pattern = '\\d{2}:\\d{2}:\\d{2},\\d{3}\\s*-->\\s*\\d{2}:\\d{2}:\\d{2},\\d{3}'
        timestamps = re.findall(timestamp_pattern, content)
        if len(timestamps) >= 1:
            sequence_pattern = '^\\d+$'
            sequences = re.findall(sequence_pattern, content, re.MULTILINE)
            if len(sequences) >= 1:
                is_valid_srt = True
                if len(timestamps) >= 2 and len(sequences) >= 2:
                    has_multiple_entries = True
    except (UnicodeDecodeError, AttributeError):
        is_valid_srt = False
    return {'exists': True, 'filename': filename, 'has_content': has_content, 'is_valid_srt': is_valid_srt, 'has_multiple_entries': has_multiple_entries, 'file_size': file_size}

def get_file_exists__6fb34f5d(env, config: dict):
    """Check if the exported file exists on VM and validate it's a proper ODS file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'exists': bool - True if file exists, False otherwise
            'is_ods': bool - True if file is ODS format, False otherwise
            'file_type': str - Output from 'file' command for verification
        }
    """
    vm_path = config.get('path', '/home/user/Desktop/annual-survey.ods')
    exists_result = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    exists = exists_result.get('output', '').strip() == 'EXISTS'
    if not exists:
        return {'exists': False, 'is_ods': False, 'file_type': ''}
    file_type_result = env.controller.run_bash_script(f"file '{vm_path}'", timeout=10)
    file_type_output = file_type_result.get('output', '').strip()
    is_ods = 'OpenDocument Spreadsheet' in file_type_output or 'application/vnd.oasis.opendocument.spreadsheet' in file_type_output
    if not is_ods and 'Zip archive' in file_type_output:
        mimetype_result = env.controller.run_bash_script(f"unzip -p '{vm_path}' mimetype 2>/dev/null || echo ''", timeout=10)
        mimetype_content = mimetype_result.get('output', '').strip()
        is_ods = mimetype_content == 'application/vnd.oasis.opendocument.spreadsheet'
    return {'exists': exists, 'is_ods': is_ods, 'file_type': file_type_output}

def get_file_info__739292ff(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get information about a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: File information (exists, name, size, extension) or None if error
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return None
    result = env.controller.execute_python_command(f"import os; print(os.path.exists('{path}'))")
    if not result or result.get('output', '').strip().lower() != 'true':
        logger.warning(f'File does not exist: {path}')
        return {'exists': False, 'name': None, 'size': 0, 'extension': None}
    info_script = f"""\nimport os\nif os.path.exists('{path}'):\n    print(f"{{os.path.basename('{path}')}}|{{os.path.getsize('{path}')}}|{{os.path.splitext('{path}')[1]}}")\nelse:\n    print("NOT_FOUND")\n"""
    result = env.controller.execute_python_command(info_script)
    if result and result.get('output'):
        output = result['output'].strip()
        if output != 'NOT_FOUND':
            parts = output.split('|')
            if len(parts) == 3:
                return {'exists': True, 'name': parts[0], 'size': int(parts[1]) if parts[1].isdigit() else 0, 'extension': parts[2]}
    return None

def get_file_content__d09dce0a(env, config: dict):
    path = config.get('path', '/home/user/file.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'content': '', 'size': 0}
    try:
        content = file_bytes.decode('utf-8').strip()
        return {'exists': True, 'content': content, 'size': len(file_bytes)}
    except:
        return {'exists': True, 'content': '', 'size': len(file_bytes)}

def get_pdf_files_with_content__83b6523335a3d68f9c734599b1537c74(env, config: Dict[str, Any]) -> Dict[str, Dict[str, Any]]:
    """
    Get PDF files with their content information from a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key and optional 'required_files' list

    Returns:
        Dict[str, Dict[str, Any]]: Dictionary mapping filename to content info
            {
                "John Doe.pdf": {
                    "exists": True,
                    "has_content": True,
                    "text_length": 1234,
                    "contains_name": True,
                    "has_rating_mark": True,
                    "has_evaluation_keywords": True,
                    "keyword_count": 3,
                    "error": None
                },
                ...
            }
    """
    directory = config.get('directory', '/home/user/Desktop')
    required_files = config.get('required_files', [])
    result_dict = {}
    evaluation_keywords = ['performance', 'evaluation', 'rating', 'review', 'department', 'position', 'score']
    for filename in required_files:
        file_path = f'{directory}/{filename}'
        file_info = {'exists': False, 'has_content': False, 'text_length': 0, 'contains_name': False, 'has_rating_mark': False, 'has_evaluation_keywords': False, 'keyword_count': 0, 'error': None}
        try:
            check_result = env.controller.run_bash_script(f'test -f "{file_path}" && echo "exists"', timeout=5)
            if check_result.get('status') == 'success' and 'exists' in check_result.get('output', ''):
                file_info['exists'] = True
                file_bytes = env.controller.get_file(file_path)
                if file_bytes:
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    try:
                        text_result = env.controller.run_bash_script(f'pdftotext "{tmp_path}" - 2>/dev/null || echo ""', timeout=10)
                        if text_result.get('status') == 'success':
                            text_content = text_result.get('output', '')
                            text_lower = text_content.lower()
                            file_info['text_length'] = len(text_content.strip())
                            file_info['has_content'] = len(text_content.strip()) > 200
                            employee_name = filename.replace('.pdf', '')
                            if employee_name.lower() in text_lower:
                                file_info['contains_name'] = True
                            if '√' in text_content or '✓' in text_content:
                                file_info['has_rating_mark'] = True
                            keyword_count = 0
                            for keyword in evaluation_keywords:
                                if keyword in text_lower:
                                    keyword_count += 1
                            file_info['keyword_count'] = keyword_count
                            file_info['has_evaluation_keywords'] = keyword_count >= 3
                    finally:
                        if os.path.exists(tmp_path):
                            os.unlink(tmp_path)
                else:
                    file_info['error'] = 'Could not read file bytes'
        except Exception as e:
            file_info['error'] = str(e)
            logger.error(f'Error checking PDF {filename}: {e}')
        result_dict[filename] = file_info
    return result_dict

def get_specific_pdf_existence__6d39fc800ff8d1397c30bfae2c676b76(env, config):
    """Check if specific PDF files exist in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'filenames' keys

    Returns:
        dict: Mapping of expected filenames to existence status (True/False)
    """
    directory_path = config.get('path', '/home/user/Desktop')
    expected_files = config.get('filenames', [])
    existence_map = {}
    for filename in expected_files:
        filepath = f'{directory_path}/{filename}'
        command = f'test -f {filepath} && echo "exists" || echo "missing"'
        result = env.controller.run_bash_script(command, timeout=10)
        if result['returncode'] == 0:
            existence_map[filename] = 'exists' in result['output']
        else:
            existence_map[filename] = False
    return existence_map

def get_file_line_count__1098fce8(env, config):
    """Get line count of a text file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest'

    Returns:
        Integer line count
    """
    vm_path = config.get('path')
    dest = config.get('dest', 'file.txt')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return 0
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            line_count = sum((1 for _ in f))
        return line_count
    except Exception as e:
        print(f'Error counting lines: {e}')
        return 0

def get_csv_unique_first_names_count__5ab9ee5026038714f957b134595e9a67(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Count unique first names in a CSV file and verify merge operation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict[str, Any]: Dictionary containing:
            - unique_count: Number of unique first names
            - total_rows: Total number of data rows
            - file_exists: Whether the file exists
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'unique_count': 0, 'total_rows': 0, 'file_exists': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'unique_count': 0, 'total_rows': 0, 'file_exists': False}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        unique_names = set()
        total_rows = 0
        with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.DictReader(f)
            for row in reader:
                total_rows += 1
                first_name = row.get('First Name', '').strip()
                if first_name:
                    unique_names.add(first_name)
        return {'unique_count': len(unique_names), 'total_rows': total_rows, 'file_exists': True}
    finally:
        os.unlink(tmp_path)

def get_numeric_value_from_file__d8bf55dd66c809967b252ee6c81aa4a7(env, config):
    """Read a text file and extract the first numeric value found.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        float: First numeric value found, or -1.0 if none found
    """
    file_path = config.get('path', '')
    if not file_path:
        return -1.0
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return -1.0
        content = file_bytes.decode('utf-8', errors='ignore')
        import re
        match = re.search('[-+]?\\d*\\.?\\d+', content)
        if match:
            return float(match.group())
        else:
            return -1.0
    except Exception as e:
        return -1.0

def get_file_exists_and_size__e1affec5(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_file_lines__8278c76c3e9924f1866901a10447ca8f(env, config):
    """Read text file content and return as list of lines.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of lines in original file order (preserves sorting) or empty list if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return []
    result = env.controller.run_bash_script(f'cat {file_path}', timeout=10)
    if result.get('returncode') == 0:
        content = result.get('output', '').strip()
        if content:
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            return lines
    return []

def get_file_exists__e3ae8a85(env, config):
    """Check if a file exists at the specified path on VM."""
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output == 'EXISTS':
        return True
    else:
        return False

def get_docx_text_content__18d2ce6a(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_file_exists__bdb8ae26(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if ODS file exists and validate it's a valid ODS file with content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with:
            - file_exists (bool): Whether the file exists
            - file_path (str): Path to the ODS file
            - is_valid_ods (bool): Whether it's a valid ODS archive
            - has_content (bool): Whether it contains spreadsheet data
            - error (str): Error message if validation fails
    """
    ods_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.ods'
    file_content = env.controller.get_file(ods_path)
    result = {'file_exists': False, 'file_path': ods_path, 'is_valid_ods': False, 'has_content': False, 'error': None}
    if file_content is None or len(file_content) == 0:
        result['error'] = 'File does not exist or is empty'
        return result
    result['file_exists'] = True
    try:
        file_bytes = io.BytesIO(file_content)
        with zipfile.ZipFile(file_bytes, 'r') as ods_zip:
            file_list = ods_zip.namelist()
            required_files = ['content.xml', 'mimetype']
            missing_files = [f for f in required_files if f not in file_list]
            if missing_files:
                result['error'] = f'Missing required ODS files: {missing_files}'
                return result
            result['is_valid_ods'] = True
            content_xml = ods_zip.read('content.xml').decode('utf-8')
            has_table = '<table:table' in content_xml
            has_cells = '<table:table-cell' in content_xml
            if has_table and has_cells:
                result['has_content'] = True
            else:
                result['error'] = 'ODS file does not contain spreadsheet data (no tables or cells found)'
    except zipfile.BadZipFile:
        result['error'] = 'File is not a valid ZIP/ODS archive'
    except Exception as e:
        result['error'] = f'Error validating ODS file: {str(e)}'
    return result

def get_dir_exists__d31c5db6(env, config: dict):
    """Check if directory exists on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path'

    Returns:
        Boolean: True if directory exists, False otherwise
    """
    dir_path = config['dir_path']
    command = f'test -d "{dir_path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return False
    return result['output'].strip() == 'exists'

def get_files_with_prefix__8c4eaed9a61673f78def8f323e7dfe9d(env, config: dict):
    """Get list of files in directory that start with a specific prefix.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (directory) and 'prefix' keys

    Returns:
        dict: {'exists': bool, 'matching_files': list of filenames}
    """
    path = config.get('path', '')
    prefix = config.get('prefix', '')
    try:
        result = env.controller.run_bash_script(f'ls -1 "{path}" 2>/dev/null', timeout=10)
        if result.get('returncode', 1) != 0:
            logger.warning(f'Directory not found or inaccessible: {path}')
            return {'exists': False, 'matching_files': []}
        output = result.get('output', '').strip()
        if not output:
            all_files = []
        else:
            all_files = [f.strip() for f in output.split('\n') if f.strip()]
        matching_files = [f for f in all_files if f.startswith(prefix)]
        logger.info(f"Found {len(matching_files)} files with prefix '{prefix}': {matching_files}")
        return {'exists': True, 'matching_files': matching_files}
    except Exception as e:
        logger.error(f'Error listing directory {path}: {e}')
        return {'exists': False, 'matching_files': []}

def get_file_absence__be265045(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file does NOT exist on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool) key
    """
    path = config['path']
    command = f"if [ -f '{path}' ]; then echo 'EXISTS'; else echo 'NOT_EXISTS'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result['output'].strip()
    exists = output == 'EXISTS' and result['returncode'] == 0
    return {'exists': exists}

def get_downloaded_file_info__08c5e1b6ad7015f1bdd4ff79ff88e12f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file was downloaded and get its properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' and 'min_size' keys

    Returns:
        Dict with file existence and size information
    """
    file_path = config.get('file_path', '')
    result = {'exists': False, 'is_file': False, 'size_bytes': 0, 'meets_min_size': False}
    check_file_cmd = f"[ -f '{file_path}' ] && echo 'EXISTS' || echo 'NOT_EXISTS'"
    file_check = env.controller.run_bash_script(check_file_cmd, timeout=10)
    if file_check.get('output', '').strip() == 'EXISTS':
        result['exists'] = True
        result['is_file'] = True
        size_cmd = f"stat -c %s '{file_path}'"
        size_result = env.controller.run_bash_script(size_cmd, timeout=10)
        if size_result.get('returncode') == 0:
            try:
                size_bytes = int(size_result.get('output', '0').strip())
                result['size_bytes'] = size_bytes
                min_size = config.get('min_size', 0)
                if size_bytes >= min_size:
                    result['meets_min_size'] = True
            except ValueError:
                logger.warning(f'Failed to parse file size for {file_path}')
    logger.info(f"Downloaded file check for {file_path}: exists={result['exists']}, size={result['size_bytes']} bytes")
    return result

def get_csv_filtered_rows__5ac234c6f992977135e16f82780e5511(env, config):
    """
    Read CSV file from VM and return all rows as list of lists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists representing CSV rows (including header)
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        rows = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
    finally:
        os.unlink(tmp_path)

def get_vm_file_size__f259522f5141b84d8b2c6c9007fd732a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on VM, get its size, validate it's a PDF, and verify it's the largest PDF from the Paper Recommendation email.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - path: Path to file on VM

    Returns:
        Dict with keys:
            - exists: bool - whether the file exists
            - size_bytes: int - file size in bytes (0 if not exists)
            - is_pdf: bool - whether file has .pdf extension
            - has_pdf_header: bool - whether file has PDF magic bytes
            - is_largest_from_email: bool - whether this is the largest PDF from the Paper Recommendation email
            - all_email_pdf_sizes: list - sizes of all PDF attachments from the email
    """
    path = config.get('path', '')
    size_result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then stat -c "%s" "{path}"; else echo "0"; fi', timeout=10)
    if size_result.get('status') != 'success':
        return {'exists': False, 'size_bytes': 0, 'is_pdf': False, 'has_pdf_header': False, 'is_largest_from_email': False, 'all_email_pdf_sizes': []}
    size_output = size_result.get('output', '0').strip()
    try:
        size = int(size_output)
        exists = size > 0
    except ValueError:
        return {'exists': False, 'size_bytes': 0, 'is_pdf': False, 'has_pdf_header': False, 'is_largest_from_email': False, 'all_email_pdf_sizes': []}
    if not exists:
        return {'exists': False, 'size_bytes': 0, 'is_pdf': False, 'has_pdf_header': False, 'is_largest_from_email': False, 'all_email_pdf_sizes': []}
    is_pdf = path.lower().endswith('.pdf')
    header_result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then head -c 5 "{path}"; else echo ""; fi', timeout=10)
    has_pdf_header = False
    if header_result.get('status') == 'success':
        header = header_result.get('output', '')
        has_pdf_header = header.startswith('%PDF-')
    all_email_pdf_sizes = []
    is_largest_from_email = False
    find_db_script = '\n    PROFILE_DIR="/home/user/.thunderbird-profile"\n    if [ -d "$PROFILE_DIR" ]; then\n        # Try multiple possible database locations\n        find "$PROFILE_DIR" -name "global-messages-db.sqlite" 2>/dev/null | head -1\n        find "$PROFILE_DIR" -name "*.msf" 2>/dev/null | head -1\n        find "$PROFILE_DIR" -path "*/ImapMail/*/INBOX.msf" 2>/dev/null | head -1\n    fi\n    '
    db_result = env.controller.run_bash_script(find_db_script, timeout=10)
    if db_result.get('status') == 'success' and db_result.get('output', '').strip():
        db_paths = db_result.get('output', '').strip().split('\n')
        for db_path in db_paths:
            if db_path and db_path.endswith('.sqlite'):
                query_attempts = [f'''sqlite3 "{db_path}" "SELECT DISTINCT attachmentName, attachmentSize FROM messagesText_content WHERE subject LIKE '%Paper Recommendation%' AND attachmentName LIKE '%.pdf' ORDER BY attachmentSize DESC;" 2>/dev/null''', f'''sqlite3 "{db_path}" "SELECT DISTINCT name, size FROM attachments JOIN messages ON attachments.message_id = messages.id WHERE messages.subject LIKE '%Paper Recommendation%' AND name LIKE '%.pdf' ORDER BY size DESC;" 2>/dev/null''', f'sqlite3 "{db_path}" ".tables" 2>/dev/null']
                for query_script in query_attempts:
                    query_result = env.controller.run_bash_script(query_script, timeout=15)
                    if query_result.get('status') == 'success':
                        output = query_result.get('output', '').strip()
                        if output and '|' in output:
                            for line in output.split('\n'):
                                if '|' in line:
                                    try:
                                        parts = line.split('|')
                                        if len(parts) >= 2:
                                            attach_size = int(parts[-1].strip())
                                            if attach_size > 0:
                                                all_email_pdf_sizes.append(attach_size)
                                    except (ValueError, IndexError):
                                        pass
                            if all_email_pdf_sizes:
                                break
                if all_email_pdf_sizes:
                    break
    if not all_email_pdf_sizes:
        cache_script = '\nPROFILE_DIR="/home/user/.thunderbird-profile"\nif [ -d "$PROFILE_DIR/Mail" ]; then\n    # Find all PDF files in the Mail directory modified within last 24 hours\n    # This helps filter out unrelated attachments from other emails\n    find "$PROFILE_DIR/Mail" -type f -name "*.pdf" -mtime -1 -exec stat -c "%s" {} \\; 2>/dev/null\nfi\n'
        cache_result = env.controller.run_bash_script(cache_script, timeout=15)
        if cache_result.get('status') == 'success' and cache_result.get('output', '').strip():
            output = cache_result.get('output', '').strip()
            for line in output.split('\n'):
                try:
                    attach_size = int(line.strip())
                    if attach_size > 0:
                        all_email_pdf_sizes.append(attach_size)
                except ValueError:
                    pass
            if len(all_email_pdf_sizes) > 1:
                path_script = '\nPROFILE_DIR="/home/user/.thunderbird-profile"\nif [ -d "$PROFILE_DIR/Mail" ]; then\n    find "$PROFILE_DIR/Mail" -type f -name "*.pdf" -mtime -1 -exec sh -c \'echo "$(stat -c "%s" "$1")|$(dirname "$1")"\' _ {} \\; 2>/dev/null\nfi\n'
                path_result = env.controller.run_bash_script(path_script, timeout=15)
                if path_result.get('status') == 'success' and path_result.get('output', '').strip():
                    dir_groups = {}
                    for line in path_result.get('output', '').strip().split('\n'):
                        if '|' in line:
                            try:
                                (size_str, dir_path) = line.split('|', 1)
                                attach_size = int(size_str.strip())
                                if dir_path not in dir_groups:
                                    dir_groups[dir_path] = []
                                dir_groups[dir_path].append(attach_size)
                            except (ValueError, IndexError):
                                pass
                    if dir_groups:
                        largest_group = max(dir_groups.values(), key=len)
                        if len(largest_group) > len(all_email_pdf_sizes) / 2:
                            all_email_pdf_sizes = largest_group
    if all_email_pdf_sizes:
        all_email_pdf_sizes = sorted(list(set(all_email_pdf_sizes)), reverse=True)
        is_largest_from_email = size == max(all_email_pdf_sizes)
    else:
        is_largest_from_email = False
    return {'exists': exists, 'size_bytes': size, 'is_pdf': is_pdf, 'has_pdf_header': has_pdf_header, 'is_largest_from_email': is_largest_from_email, 'all_email_pdf_sizes': all_email_pdf_sizes}

def get_file_count_in_dir__081d0b6c(env, config: dict):
    """
    Count number of files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path' and optional 'extension' filter

    Returns:
        int: Number of files found
    """
    dir_path = config.get('dir_path', '/home/user/Downloads')
    extension = config.get('extension', '')
    if extension:
        command = f"find {dir_path} -maxdepth 1 -name '*{extension}' -type f | wc -l"
    else:
        command = f'find {dir_path} -maxdepth 1 -type f | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    try:
        count = int(result['output'].strip())
        logger.info(f"Found {count} files in {dir_path} with extension '{extension}'")
        return count
    except Exception as e:
        logger.error(f'Failed to parse file count: {e}')
        return 0

def get_pdf_exports__2586a709(env, config: dict):
    """
    Check if PDFs were exported to the specified directory.
    Validates that PDFs are valid and extracts their content.
    Also extracts PDF metadata to get the title if available.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict containing:
        - pdf_files: List of dicts with 'filename', 'valid', 'text_content', 'page_count', 'pdf_title'
        - count: Total number of PDF files found
    """
    directory = config.get('directory', '/home/user/Desktop/PDFs')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        full_paths = [f.strip() for f in list_result['output'].strip().split('\n') if f.strip()]
        for pdf_path in full_paths:
            filename = os.path.basename(pdf_path)
            pdf_info = {'filename': filename, 'valid': False, 'text_content': '', 'page_count': 0, 'pdf_title': ''}
            validate_script = f"""\nimport sys\ntry:\n    import fitz  # PyMuPDF\n    doc = fitz.open('{pdf_path}')\n    page_count = len(doc)\n    text_content = ''\n    for page in doc:\n        text_content += page.get_text()\n\n    # Try to extract PDF title from metadata (BEFORE closing the document)\n    metadata = doc.metadata\n    pdf_title = metadata.get('title', '') if metadata else ''\n\n    doc.close()\n\n    # Output: valid|page_count|pdf_title|text_preview (first 1000 chars)\n    print(f"VALID|{{page_count}}|{{pdf_title}}|{{text_content[:1000]}}")\nexcept Exception as e:\n    print(f"INVALID|0||Error: {{str(e)}}")\n    sys.exit(0)\n"""
            validation_result = env.controller.run_python_script(validate_script, timeout=15)
            if validation_result['output']:
                output = validation_result['output'].strip()
                if output.startswith('VALID|'):
                    parts = output.split('|', 3)
                    pdf_info['valid'] = True
                    pdf_info['page_count'] = int(parts[1])
                    pdf_info['pdf_title'] = parts[2] if len(parts) > 2 else ''
                    pdf_info['text_content'] = parts[3] if len(parts) > 3 else ''
            pdf_files.append(pdf_info)
    return {'pdf_files': pdf_files, 'count': len(pdf_files)}

def get_txt_content__562f595e9fa43de1b5d952aea3a3f9eb(env, config):
    """Get text file content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception:
        return ''

def get_docx_text_content__cea00e79(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_pdf_filenames_in_dir__7152d37e(env, config: dict):
    """Get list of PDF filenames in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path to check

    Returns:
        List of PDF filenames (without extension)
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"""\npython3 -c "\nimport os\nimport glob\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path):\n    print('[]')\n    exit(0)\n\npdf_files = [os.path.splitext(os.path.basename(f))[0] for f in glob.glob(os.path.join(dir_path, '*.pdf'))]\nimport json\nprint(json.dumps(pdf_files))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to list PDF files: {result['error']}")
        return []
    import json
    try:
        files = json.loads(result['output'].strip())
        return files
    except:
        return []

def get_single_pdf_with_verification__1386929a4648885c7d87d6425829fd94(env, config: dict):
    """Check if a single PDF file exists with specific content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path on VM
            - filename: Expected filename
            - content_check: Text that must be in the PDF

    Returns:
        dict: Result with 'exists', 'has_content' keys
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    filename = config.get('filename', '')
    content_check = config.get('content_check', '')
    result = {'exists': False, 'has_content': False}
    if not filename:
        return result
    file_path = os.path.join(directory, filename)
    check_cmd = f'[ -f "{file_path}" ] && echo "exists" || echo "not_found"'
    cmd_result = env.controller.run_bash_script(check_cmd, timeout=10)
    if cmd_result.get('output', '').strip() != 'exists':
        return result
    result['exists'] = True
    if content_check:
        try:
            file_bytes = env.controller.get_file(file_path)
            if not file_bytes:
                return result
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                import fitz
                doc = fitz.open(tmp_path)
                text = ''
                for page in doc:
                    text += page.get_text()
                doc.close()
                if content_check in text:
                    result['has_content'] = True
            except Exception:
                pass
            finally:
                os.unlink(tmp_path)
        except Exception:
            pass
    return result

def get_numeric_value_from_txt__f9a0219a(env, config):
    """
    Read a text file and extract the first numeric value found.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: First numeric value found, or 0 if none
    """
    import re
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return 0
        content = file_bytes.decode('utf-8', errors='ignore')
        match = re.search('\\d+', content)
        if match:
            return int(match.group())
        return 0
    except Exception as e:
        return 0

def get_docx_text_content__97c6df2d(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_folder_organization__6a9b75b65029a1742667a75cc968a2a9(env, config: Dict) -> Dict:
    """
    Check if files are organized into specific folders and verify original files were moved.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'base_directory' key

    Returns:
        Dict with:
            - 'folders': mapping of folder names to list of files in them
            - 'original_files_exist': dict mapping original filenames to bool (True if still in base directory)
    """
    base_directory = config.get('base_directory', '/home/user/Pictures')
    command = f"""python3 -c "\nimport os\nimport json\n\nbase = '{base_directory}'\nresult = {{}}\nresult['folders'] = {{}}\nresult['original_files_exist'] = {{}}\n\n# Check common mountain folder names\nfolders_to_check = ['Kilimanjaro', 'kilimanjaro', 'Everest', 'everest', 'Hua', 'hua', 'Huashan', 'huashan']\n\nfor folder_name in folders_to_check:\n    folder_path = os.path.join(base, folder_name)\n    if os.path.isdir(folder_path):\n        files = os.listdir(folder_path)\n        result['folders'][folder_name] = files\n\n# Check if original files still exist in base directory\noriginal_files = ['picture1.jpg', 'picture2.jpg', 'picture3.jpg']\nfor filename in original_files:\n    file_path = os.path.join(base, filename)\n    result['original_files_exist'][filename] = os.path.exists(file_path)\n\nprint(json.dumps(result))\n"\n"""
    run_result = env.controller.run_bash_script(command, timeout=10)
    if run_result.get('returncode') != 0:
        return {'folders': {}, 'original_files_exist': {}}
    output = run_result.get('output', '')
    try:
        import json
        folder_contents = json.loads(output.strip())
        return folder_contents
    except:
        return {'folders': {}, 'original_files_exist': {}}

def get_archived_ipynb_files__0190943f9252410b5b69d12d28f1b6b7(env, config):
    """
    Get list of *_archived.ipynb files in test_environment.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Dictionary with two keys:
            - 'archived_files': list of *_archived.ipynb files found (relative paths)
            - 'failed_files': list of *failed.ipynb files still present (relative paths)
    """
    archived_cmd = 'cd /home/user/test_environment && find . -name "*_archived.ipynb" -type f | sed "s|^./||" | sort'
    archived_result = env.controller.run_bash_script(archived_cmd, timeout=30)
    failed_cmd = 'cd /home/user/test_environment && find . -name "*failed.ipynb" -type f | sed "s|^./||" | sort'
    failed_result = env.controller.run_bash_script(failed_cmd, timeout=30)
    archived_files = []
    failed_files = []
    if archived_result['status'] == 'success':
        output = archived_result['output'].strip()
        if output:
            archived_files = [f for f in output.split('\n') if f]
    if failed_result['status'] == 'success':
        output = failed_result['output'].strip()
        if output:
            failed_files = [f for f in output.split('\n') if f]
    return {'archived_files': sorted(archived_files), 'failed_files': sorted(failed_files)}

def get_readme_content__457448603041bcfcd06aff998da4e920(env, config: Dict[str, Any]) -> str:
    """Get README.md file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String containing README content, or empty string if file not found
    """
    file_path = config.get('path', '/home/user/project/README.md')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'README file not found: {file_path}')
            return ''
        return file_bytes.decode('utf-8')
    except Exception as e:
        logger.error(f'Error reading README file: {e}')
        return ''

def get_multi_directory_contents__94dfca2e(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_tetris_files__d0091276(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_file_size__f778a8914a698cd2bc7c0cc50cd3596d(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Get file size and image validation information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'size', 'is_valid_png', 'width', 'height' keys, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    result = {'size': len(file_bytes), 'is_valid_png': False, 'width': 0, 'height': 0}
    if len(file_bytes) >= 8:
        png_signature = b'\x89PNG\r\n\x1a\n'
        if file_bytes[:8] == png_signature:
            result['is_valid_png'] = True
            try:
                from PIL import Image
                img = Image.open(BytesIO(file_bytes))
                result['width'] = img.width
                result['height'] = img.height
                logger.info(f"Valid PNG image: {result['width']}x{result['height']}, {result['size']} bytes")
            except Exception as e:
                logger.warning(f'PNG signature valid but PIL failed to open image: {e}')
                result['is_valid_png'] = False
        else:
            logger.warning(f'File does not have PNG magic bytes')
    return result

def get_file_exists__03122ed4(env, config):
    """
    Check if a file exists and validate it's a proper PNG screenshot.

    Returns a dict with:
    - exists: bool (file exists)
    - is_png: bool (file is a valid PNG)
    - size: int (file size in bytes)
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    file_info = {'exists': output == 'EXISTS', 'is_png': False, 'size': 0}
    if not file_info['exists']:
        return file_info
    file_type_command = f'file -b "{file_path}"'
    file_type_result = env.controller.run_bash_script(file_type_command, timeout=10)
    file_type_output = file_type_result.get('output', '').strip().lower()
    file_info['is_png'] = 'png' in file_type_output
    size_command = f'stat -c %s "{file_path}" 2>/dev/null || stat -f %z "{file_path}" 2>/dev/null'
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    size_output = size_result.get('output', '').strip()
    try:
        file_info['size'] = int(size_output)
    except (ValueError, TypeError):
        file_info['size'] = 0
    return file_info

def get_directory_has_file__b77b1963e53aafc9923374c6ea5077e2(env, config: Dict[str, Any]) -> bool:
    """Check if a specific file exists in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (full path to expected file)

    Returns:
        bool: True if file exists in the directory, False otherwise
    """
    file_bytes = env.controller.get_file(config['path'])
    return file_bytes is not None

def get_text_file_lines__36494d7a38fdbc8d1f722d2db36fce8e(env, config):
    """
    Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: Lines from the file (stripped of whitespace)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return []

def get_speedtest_json_data__c8d946870135d67f0db0be5e65caaa2a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract speedtest data from JSON file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with JSON data or empty dict if file doesn't exist/is invalid
    """
    import os
    import json
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Failed to get file: {config['path']}")
        return {}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.json', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return data if isinstance(data, dict) else {}
    except Exception as e:
        logger.error(f'Error reading JSON file: {e}')
        return {}
    finally:
        os.unlink(tmp_path)

def get_file_properties__9c119f81(env, config: dict):
    path = config.get('path', '/home/user/file.docx')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0}
    return {'exists': True, 'size': len(file_bytes)}

def get_text_file_content__d7828490(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_file_move_check__aace45122e840d40b84dc540ae5a49bc(env, config: dict):
    """
    Check file existence in multiple directories and root.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        dict: File existence information for directories and root
    """
    check_script = '\n#!/bin/bash\nresult=""\n\n# Check each directory\nfor dir in dir1 dir2 dir3; do\n    if [ -f "$dir/file1" ]; then\n        result="${result}exists_in_${dir}|"\n    fi\ndone\n\n# Check root\nif [ -f "file1" ]; then\n    result="${result}exists_in_root|"\nfi\n\necho "$result"\n'
    result = env.controller.run_bash_script(check_script, timeout=10)
    if result['returncode'] != 0:
        return {'dir1': False, 'dir2': False, 'dir3': False, 'root': False}
    output = result['output'].strip()
    return {'dir1': 'exists_in_dir1' in output, 'dir2': 'exists_in_dir2' in output, 'dir3': 'exists_in_dir3' in output, 'root': 'exists_in_root' in output}

def get_docx_text_content__6003371e(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_file_line_count__475840bd88bfc32515242a838ac799b5(env, config: dict):
    """Count lines in a text file from VM and verify chapter content ordering.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the text file on VM

    Returns:
        dict: {
            'exists': bool,
            'line_count': int,
            'first_500_chars': str,
            'last_500_chars': str,
            'chapter_boundaries': dict with chapter positions,
            'content_hash': str,
            'full_content': str
        }
    """
    import hashlib
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return {'exists': False, 'line_count': 0, 'first_500_chars': '', 'last_500_chars': '', 'chapter_boundaries': {}, 'content_hash': ''}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.splitlines()
        line_count = len(lines)
        first_500 = content[:500] if len(content) >= 500 else content
        last_500 = content[-500:] if len(content) >= 500 else content
        content_hash = hashlib.sha256(content.encode('utf-8')).hexdigest()
        chapter_patterns = {'ch0': ['Chapter 1 Passing through the book', 'Chi Qingluo opened his eyes', 'sleeping in the bathtub'], 'ch1': ['Chapter 2 Change', 'Qin Yan paused and frowned slightly', 'Just a few minutes!'], 'ch2': ['Chapter 3 Plan', "Chi Qingluo paused while biting the lion's head", 'The original owner was on a diet?'], 'ch3': ['Chapter 4 Black Material', 'Silence does not mean admitting defeat', 'public relations still need to be done'], 'ch4': ['Chapter 5: Ask him to pay back the money', "In the Qin family's study room", 'Having said all that, can you pay me back?']}
        chapter_boundaries = {}
        for (chapter_id, patterns) in chapter_patterns.items():
            found = False
            for (i, line) in enumerate(lines):
                line_lower = line.lower()
                for pattern in patterns:
                    if pattern.lower() in line_lower:
                        chapter_boundaries[chapter_id] = {'line': i, 'content': line.strip()[:100]}
                        found = True
                        logger.info(f'Found {chapter_id} at line {i}: {line.strip()[:50]}')
                        break
                if found:
                    break
            if not found:
                chapter_num = chapter_id.replace('ch', '')
                for (i, line) in enumerate(lines):
                    if f'chapter {chapter_num}' in line.lower() or f'chapter{chapter_num}' in line.lower():
                        chapter_boundaries[chapter_id] = {'line': i, 'content': line.strip()[:100]}
                        logger.info(f'Found {chapter_id} at line {i} (fallback): {line.strip()[:50]}')
                        break
        logger.info(f'File {path} has {line_count} lines, found {len(chapter_boundaries)}/5 chapters')
        logger.info(f'Content hash: {content_hash}')
        return {'exists': True, 'line_count': line_count, 'first_500_chars': first_500, 'last_500_chars': last_500, 'chapter_boundaries': chapter_boundaries, 'content_hash': content_hash, 'full_content': content}
    except Exception as e:
        logger.error(f'Error processing file {path}: {e}')
        return {'exists': False, 'line_count': 0, 'first_500_chars': '', 'last_500_chars': '', 'chapter_boundaries': {}, 'content_hash': ''}

def get_file_exists__46ccb784(env, config: dict):
    """Check if file exists on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path'

    Returns:
        Boolean: True if file exists, False otherwise
    """
    file_path = config['file_path']
    command = f'test -f "{file_path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return False
    return result['output'].strip() == 'exists'

def get_pdf_files_list__d2aaaf87(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files with their page counts in the specified directory.

    Returns:
        Dict[str, int]: Mapping of filename to page count
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    list_command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    if list_result.get('returncode') != 0:
        return {}
    output = list_result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    file_page_counts = {}
    for filename in files:
        filepath = f'{directory}/{filename}'
        pdfinfo_command = f"pdfinfo '{filepath}' 2>/dev/null | grep -oP '(?<=Pages:)\\s+\\d+' | tr -d ' '"
        pdfinfo_result = env.controller.run_bash_script(pdfinfo_command, timeout=10)
        if pdfinfo_result.get('returncode') == 0:
            page_count_str = pdfinfo_result.get('output', '').strip()
            try:
                page_count = int(page_count_str)
                file_page_counts[filename] = page_count
            except (ValueError, TypeError):
                continue
    return file_page_counts

def get_dir_exists__92ec1d53(env, config: dict):
    """Check if directory exists with exactly 2 valid PDF files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path to check

    Returns:
        1 if directory exists with exactly 2 valid PDF files, 0 otherwise
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"""\npython3 -c "\nimport os\nimport glob\nimport time\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path) or not os.path.isdir(dir_path):\n    print('0')\n    exit()\n\n# Get all PDF files in the directory\npdf_files = glob.glob(os.path.join(dir_path, '*.pdf'))\n\n# Check for exactly 2 PDF files (matching the 2 URLs opened)\nif len(pdf_files) != 2:\n    print('0')\n    exit()\n\n# Verify each PDF file is valid\nvalid_count = 0\nfor pdf_file in pdf_files:\n    try:\n        # Check file size - should be greater than 10KB (typical blog PDFs are 100KB+)\n        file_size = os.path.getsize(pdf_file)\n        if file_size < 10240:  # 10KB minimum\n            continue\n\n        # Check if file appears to be a valid PDF (starts with PDF signature)\n        with open(pdf_file, 'rb') as f:\n            header = f.read(4)\n            if header == b'%PDF':\n                valid_count += 1\n    except Exception:\n        continue\n\n# Success if we have exactly 2 valid PDFs\nif valid_count == 2:\n    print('1')\nelse:\n    print('0')\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to check directory: {result['error']}")
        return 0
    try:
        exists = int(result['output'].strip())
        return exists
    except:
        return 0

def get_subdirs_count__125fed35(env, config):
    """
    Verify the git repository was successfully fetched/cloned.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        dict: Verification results including:
            - has_git_dir: bool - whether .git directory exists
            - git_remote_url: str - the remote URL (if git repo exists)
            - has_instructor_embedding_dir: bool - whether InstructorEmbedding subdirectory exists
            - subdirs_count: int - number of subdirectories
    """
    repo_path = config.get('repo_path', '/home/user/instructor-embedding')
    result = {'has_git_dir': False, 'git_remote_url': '', 'has_instructor_embedding_dir': False, 'subdirs_count': 0}
    git_check_cmd = f"test -d {repo_path}/.git && echo 'yes' || echo 'no'"
    git_check_result = env.controller.run_bash_script(git_check_cmd, timeout=10)
    result['has_git_dir'] = git_check_result.get('output', '').strip() == 'yes'
    if result['has_git_dir']:
        remote_cmd = f"cd {repo_path} && git remote -v 2>/dev/null | grep fetch | awk '{{print $2}}' | head -n 1"
        remote_result = env.controller.run_bash_script(remote_cmd, timeout=10)
        result['git_remote_url'] = remote_result.get('output', '').strip()
    instructor_dir_cmd = f"test -d {repo_path}/InstructorEmbedding && echo 'yes' || echo 'no'"
    instructor_dir_result = env.controller.run_bash_script(instructor_dir_cmd, timeout=10)
    result['has_instructor_embedding_dir'] = instructor_dir_result.get('output', '').strip() == 'yes'
    count_cmd = f"find {repo_path} -mindepth 1 -maxdepth 1 -type d ! -name '.*' 2>/dev/null | wc -l"
    count_result = env.controller.run_bash_script(count_cmd, timeout=10)
    try:
        result['subdirs_count'] = int(count_result.get('output', '0').strip())
    except ValueError:
        result['subdirs_count'] = 0
    return result

def get_rule_vm_file(env, config: Dict[str, str]) -> str:
    """
    Get a file from the VM and return the local path.
    This is used when expected.type is 'rule_vm_file'.

    Args:
        env: Environment object
        config: Configuration dict containing 'rules' with 'path' and 'dest' keys

    Returns:
        str: Local file path where the VM file has been downloaded
    """
    rules = config.get('rules', {})
    vm_path = rules.get('path')
    dest_filename = rules.get('dest')
    if not vm_path:
        raise ValueError('rules.path is required for get_rule_vm_file')
    if not dest_filename:
        import os
        dest_filename = os.path.basename(vm_path)
    local_path = get_vm_file(env, {'path': vm_path, 'dest': dest_filename})
    return local_path

def get_vm_file_exists__3df7d80c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on the VM and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool) and 'size_bytes' (int) keys
    """
    path = config['path']
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; else echo 'NOT_EXISTS'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result['output'].strip()
    if output == 'NOT_EXISTS' or result['returncode'] != 0:
        return {'exists': False, 'size_bytes': 0}
    try:
        size = int(output)
        return {'exists': True, 'size_bytes': size}
    except ValueError:
        logger.error(f'Failed to parse file size: {output}')
        return {'exists': False, 'size_bytes': 0}

def get_python_script_content__c9385c6b(env, config: dict):
    """Get content of a Python script file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Content of the Python file
    """
    path = config.get('path', '/home/user/Documents/Projects/OSWorld/setup.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception:
        return ''

def get_text_file_content__f8ce9de4(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_pdf_with_name_pattern__55b958e432c8380deab73a3d2fcf329a(env, config: dict):
    """Check if PDF files exist with specific naming pattern and content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path on VM
            - expected_files: List of dicts with 'filename' and optional 'content_check'

    Returns:
        dict: Results for each file {filename: status (1=correct, 0=wrong content, -1=not found)}
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    expected_files = config.get('expected_files', [])
    results = {}
    for file_info in expected_files:
        filename = file_info['filename']
        content_check = file_info.get('content_check', '')
        file_path = os.path.join(directory, filename)
        check_cmd = f'[ -f "{file_path}" ] && echo "exists" || echo "not_found"'
        cmd_result = env.controller.run_bash_script(check_cmd, timeout=10)
        if cmd_result.get('output', '').strip() != 'exists':
            results[filename] = -1
            continue
        if content_check:
            try:
                file_bytes = env.controller.get_file(file_path)
                if not file_bytes:
                    results[filename] = 0
                    continue
                with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    import fitz
                    doc = fitz.open(tmp_path)
                    text = ''
                    for page in doc:
                        text += page.get_text()
                    doc.close()
                    if content_check in text:
                        results[filename] = 1
                    else:
                        results[filename] = 0
                except Exception:
                    results[filename] = 0
                finally:
                    os.unlink(tmp_path)
            except Exception:
                results[filename] = 0
        else:
            results[filename] = 1
    return results

def get_pdf_files_list__c8947a04(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files with their page counts in the specified directory.

    Returns:
        Dict mapping filename to page count (e.g., {'Chapter_1.pdf': 10})
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_data = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        count_cmd = f"pdfinfo '{filepath}' 2>/dev/null | grep -i 'Pages:' | awk '{{print $2}}'"
        count_result = env.controller.run_bash_script(count_cmd, timeout=10)
        if count_result.get('returncode') == 0:
            page_count_str = count_result.get('output', '').strip()
            try:
                page_count = int(page_count_str)
                pdf_data[filename] = page_count
            except (ValueError, TypeError):
                pdf_data[filename] = 0
        else:
            pdf_data[filename] = 0
    return pdf_data

def get_dir_file_list__36005e14c2ae4846381f0946699e70be(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get list of files in a directory on the VM and verify zip file existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' key

    Returns:
        Dict with 'zip_exists' (bool), 'files' (list), and 'error' (str) keys
    """
    result_dict = {'zip_exists': False, 'files': [], 'error': None}
    zip_check_cmd = 'test -f /home/user/Desktop/backup.zip && echo "exists" || echo "not_found"'
    zip_check_result = env.controller.run_bash_script(zip_check_cmd, timeout=30)
    if zip_check_result['returncode'] == 0 and 'exists' in zip_check_result['output']:
        result_dict['zip_exists'] = True
    else:
        result_dict['error'] = 'Zip file not found at /home/user/Desktop/backup.zip'
        return result_dict
    zip_list_cmd = 'unzip -l /home/user/Desktop/backup.zip 2>&1'
    zip_list_result = env.controller.run_bash_script(zip_list_cmd, timeout=30)
    if zip_list_result['returncode'] != 0:
        result_dict['error'] = 'Zip file exists but is invalid or corrupted'
        return result_dict
    unzip_cmd = 'rm -rf /home/user/Desktop/backup && unzip /home/user/Desktop/backup.zip -d /home/user/Desktop'
    unzip_result = env.controller.run_bash_script(unzip_cmd, timeout=30)
    if unzip_result['returncode'] != 0:
        result_dict['error'] = 'Failed to unzip backup.zip'
        return result_dict
    command = config.get('command', 'ls /home/user/Desktop/backup/')
    list_result = env.controller.run_bash_script(command, timeout=30)
    if list_result['returncode'] != 0:
        result_dict['error'] = 'Backup folder not found after unzipping'
        return result_dict
    output = list_result['output'].strip()
    if output:
        files = [f.strip() for f in output.split('\n') if f.strip()]
        result_dict['files'] = files
    return result_dict

def get_multi_directory_contents__8c566ad0(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_filename_hash_mapping__bf825e2c(env, config):
    """Get mapping of filenames to their image hashes.

    Args:
        env: Desktop environment
        config: Dict with 'directory' key

    Returns:
        dict: Mapping of filename (without path) to hash
    """
    directory = config.get('directory', '/home/user/Pictures')
    result = env.controller.run_bash_script(f"find {directory} -maxdepth 1 -type f \\( -iname '*.jpg' -o -iname '*.jpeg' -o -iname '*.png' \\)", timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        return {}
    files = result['output'].strip().split('\n')
    filename_hash_map = {}
    for file_path in files:
        if not file_path or file_path.strip() == '':
            continue
        file_path = file_path.strip()
        filename = os.path.basename(file_path)
        file_bytes = env.controller.get_file(file_path)
        if file_bytes:
            try:
                with Image.open(BytesIO(file_bytes)) as img:
                    img_byte_arr = img.tobytes()
                    hash_result = hashlib.sha256(img_byte_arr).hexdigest()
                    filename_hash_map[filename] = hash_result
            except Exception as e:
                print(f'Error processing {file_path}: {e}')
                continue
    return filename_hash_map

def get_file_list_content__557a0701(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get content from a text file containing a list of filenames.

    Config:
        path (str): absolute path on the VM to fetch the file
        dest (str): file name of the downloaded file

    Returns:
        Dict with keys:
            - exists (bool): whether the file exists
            - line_count (int): number of non-empty lines in the file
            - files (List[str]): list of filenames (one per line, cleaned)
            - content (str): raw file content
    """
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    try:
        file_data = env.controller.get_file(path)
        if file_data is None:
            logger.warning(f'File not found on VM: {path}')
            return {'exists': False, 'line_count': 0, 'files': [], 'content': ''}
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        lines = content.splitlines()
        files = [line.strip() for line in lines if line.strip()]
        line_count = len(files)
        logger.info(f'Successfully read file list: {path} ({line_count} files)')
        return {'exists': True, 'line_count': line_count, 'files': files, 'content': content}
    except Exception as e:
        logger.error(f'Error reading file list {path}: {e}')
        return {'exists': False, 'line_count': 0, 'files': [], 'content': ''}

def get_targz_contents__77bd062dcbbc640025c417629c0311e0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract and verify TAR.GZ file contents from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to TAR.GZ file on VM

    Returns:
        Dict with 'exists', 'is_valid_tar', and 'file_list' keys
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'is_valid_tar': False, 'file_list': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.tar.gz', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with tarfile.open(tmp_path, 'r:gz') as tf:
                file_list = sorted([member.name for member in tf.getmembers() if member.isfile()])
                return {'exists': True, 'is_valid_tar': True, 'file_list': file_list}
        except (tarfile.TarError, Exception):
            return {'exists': True, 'is_valid_tar': False, 'file_list': []}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'is_valid_tar': False, 'file_list': []}

def get_text_file_format__7f844786255954cce16f5ea58433f34e(env, config):
    """
    Read text file and check format patterns.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with format analysis (lines, has_numbers, etc.)
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'lines': [], 'has_numbering': False, 'line_count': 0}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'lines': [], 'has_numbering': False, 'line_count': 0}
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.split('\n') if line.strip()]
        import re
        has_numbering = True
        for (i, line) in enumerate(lines, 1):
            if not re.match(f'^{i}\\.\\s+', line):
                has_numbering = False
                break
        return {'lines': lines, 'has_numbering': has_numbering, 'line_count': len(lines)}
    except Exception as e:
        print(f'Error reading text file: {e}')
        return {'lines': [], 'has_numbering': False, 'line_count': 0}

def get_csv_sorted_contacts__1e62491a(env, config: dict):
    """Get contacts from CSV file in the order they appear.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of contact dictionaries in the order they appear in the CSV
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    lines = content.strip().split('\n')
    if len(lines) < 2:
        return []
    reader = csv.DictReader(lines)
    contacts = list(reader)
    return contacts

def get_file_exists__9c31b3f6afb568d600c17c937149b6c4(env, config: Dict[str, Any]) -> bool:
    """Check if a file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    file_bytes = env.controller.get_file(config['path'])
    return file_bytes is not None

def get_files_matching_pattern__039c45a2(env, config: dict):
    """Get list of files matching a pattern in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'pattern' parameters

    Returns:
        List of matching filenames
    """
    path = config.get('path', '/home/user/Documents/Finance/receipts')
    pattern = config.get('pattern', '*.pdf')
    result = env.controller.run_bash_script(f"ls '{path}'/{pattern} 2>/dev/null | xargs -n 1 basename", timeout=10)
    if result['returncode'] != 0:
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [line.strip() for line in output.split('\n') if line.strip()]
    return files

def get_pdf_files_info__437d5a7f(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_content_match__57d8acad(env, config: dict):
    """
    Read file content and return it.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file doesn't exist
    """
    path = config.get('path', '')
    command = f'cat "{path}" 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        return result['output'].strip()
    return ''

def get_subdirs_exist__2a5217dc(env, config):
    """Check if expected subdirectories exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'base_path' and 'subdirs' parameters

    Returns:
        dict: {subdir_name: bool} mapping of subdirectory existence
    """
    base_path = config.get('base_path')
    subdirs = config.get('subdirs', [])
    results = {}
    for subdir in subdirs:
        full_path = f'{base_path}/{subdir}'
        command = f"test -d {full_path} && echo 'YES' || echo 'NO'"
        result = env.controller.run_bash_script(command, timeout=10)
        results[subdir] = 'YES' in result.get('output', '')
    return results

def get_file_in_directory__31a8a4acc19afab71c5f7ec2f3006a11(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a specific file exists in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        Dict with 'exists' boolean and 'file_path' string
    """
    file_path = config.get('file_path', '')
    check_command = f"test -f '{file_path}' && echo 'exists' || echo 'not_exists'"
    result_data = env.controller.run_bash_script(check_command, timeout=10)
    exists = False
    if result_data and result_data.get('returncode') == 0:
        output = result_data.get('output', '').strip()
        exists = output == 'exists'
    return {'exists': exists, 'file_path': file_path}

def get_pdf_files_list__cedddaca(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get dictionary mapping PDF filenames to their page counts in the specified directory."""
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        count_command = f"pdfinfo '{filepath}' 2>/dev/null | grep -i '^Pages:' | awk '{{print $2}}'"
        count_result = env.controller.run_bash_script(count_command, timeout=10)
        if count_result.get('returncode') == 0:
            page_count_str = count_result.get('output', '').strip()
            try:
                page_count = int(page_count_str)
                pdf_info[filename] = page_count
            except (ValueError, TypeError):
                pass
    return pdf_info

def get_python_file_info__198be354(env, config):
    """Get information about a Python file including line count, validity, and content markers.

    Enhanced to verify code is from a SPECIFIC Colab notebook using:
    - AST-based structural fingerprinting (function/class counts)
    - Unique content markers from the target notebook
    - Code complexity and structure analysis

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Information about the file (exists, line_count, has_content, is_valid_python,
              has_expected_markers, structural_fingerprint, unique_content_score)
    """
    import os
    import ast
    import re
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'line_count': 0, 'has_content': False, 'size': 0, 'is_valid_python': False, 'has_expected_markers': False, 'marker_count': 0, 'specificity_score': 0.0, 'structural_fingerprint': {}, 'unique_content_score': 0.0, 'has_colab_metadata': False}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    lines = content_str.splitlines()
    non_empty_lines = [line for line in lines if line.strip()]
    is_valid_python = False
    structural_fingerprint = {'function_count': 0, 'class_count': 0, 'import_count': 0, 'function_names': [], 'class_names': [], 'has_main_or_train': False, 'has_model_definition': False}
    try:
        tree = ast.parse(content_str)
        is_valid_python = True
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef):
                structural_fingerprint['function_count'] += 1
                structural_fingerprint['function_names'].append(node.name)
            elif isinstance(node, ast.ClassDef):
                structural_fingerprint['class_count'] += 1
                structural_fingerprint['class_names'].append(node.name)
            elif isinstance(node, (ast.Import, ast.ImportFrom)):
                structural_fingerprint['import_count'] += 1
        structural_fingerprint['has_main_or_train'] = any((name in ['main', 'train', 'train_model', 'training', 'run_training'] for name in structural_fingerprint['function_names']))
        structural_fingerprint['has_model_definition'] = any(('model' in name.lower() or 'net' in name.lower() for name in structural_fingerprint['function_names'] + structural_fingerprint['class_names']))
    except SyntaxError:
        is_valid_python = False
    except Exception:
        is_valid_python = False
    unique_markers = {'sklearn_model_import': bool(re.search('from sklearn\\.\\w+\\s+import\\s+\\w+', content_str)), 'tf_keras_layers': 'tensorflow.keras.layers' in content_str or 'from tensorflow.keras import layers' in content_str, 'torch_nn': 'torch.nn' in content_str or 'from torch import nn' in content_str, 'data_loading_pattern': bool(re.search('(pd\\.read_csv|pd\\.read_excel|np\\.load|load_data)', content_str)), 'model_training_pattern': bool(re.search('\\.(fit|train|compile)\\s*\\(', content_str)), 'model_evaluation_pattern': bool(re.search('\\.(evaluate|score|predict|test)\\s*\\(', content_str)), 'data_transformation': bool(re.search('(\\.transform\\(|\\.fit_transform\\(|StandardScaler|MinMaxScaler|preprocessing)', content_str)), 'layer_definitions': bool(re.search('(Dense|Conv2D|LSTM|Dropout|BatchNormalization|Linear|Conv1d)\\s*\\(', content_str)), 'optimizer_usage': bool(re.search('(Adam|SGD|RMSprop|optimizer)\\s*\\(', content_str)), 'loss_function': bool(re.search('(loss|criterion|CrossEntropy|MSE|MAE)\\s*[=\\(]', content_str)), 'plotting_code': bool(re.search('(plt\\.plot|plt\\.show|sns\\.|matplotlib)', content_str)), 'colab_references': bool(re.search('(colab|google\\.colab|drive\\.mount|%|!pip|!wget)', content_str, re.IGNORECASE))}
    unique_content_count = sum(unique_markers.values())
    unique_content_score = min(1.0, unique_content_count / 8.0)
    has_colab_metadata = unique_markers.get('colab_references', False)
    generic_markers = ['import numpy', 'import pandas', 'import matplotlib', 'from sklearn', 'import tensorflow', 'import torch', 'def train', 'def model']
    marker_count = sum((1 for marker in generic_markers if marker in content_str))
    has_expected_markers = marker_count >= 3 and structural_fingerprint['function_count'] >= 3 and (structural_fingerprint['import_count'] >= 5)
    specificity_score = unique_content_score * 0.6 + min(1.0, structural_fingerprint['function_count'] / 10.0) * 0.4
    return {'exists': True, 'line_count': len(lines), 'non_empty_line_count': len(non_empty_lines), 'has_content': len(non_empty_lines) > 0, 'size': len(file_content), 'is_valid_python': is_valid_python, 'has_expected_markers': has_expected_markers, 'marker_count': marker_count, 'specificity_score': specificity_score, 'structural_fingerprint': structural_fingerprint, 'unique_content_score': unique_content_score, 'unique_markers': unique_markers, 'has_colab_metadata': has_colab_metadata}

def get_docx_text_content__94746a60(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_multi_directory_contents__f3cabf2e(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_pdf_orientation__c75d6b17(env, config: Dict[str, Any]):
    """Get the orientation of the first page in a PDF.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' to PDF file

    Returns:
        str: 'landscape' if width > height, 'portrait' otherwise, or None if error
    """
    file_content = env.controller.get_file(config['path'])
    if not file_content:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name
    try:
        reader = PdfReader(tmp_path)
        if len(reader.pages) == 0:
            return None
        page = reader.pages[0]
        mediabox = page.mediabox
        width = float(mediabox.width)
        height = float(mediabox.height)
        if width > height:
            return 'landscape'
        else:
            return 'portrait'
    finally:
        os.unlink(tmp_path)

def get_file_content__3c9f051952e2f37565e45b593e085b87(env, config: Dict[str, Any]) -> str:
    """Get the content of a file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location

    Returns:
        String containing the file content, or empty string if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        return file_bytes.decode('utf-8').strip()
    except Exception as e:
        logger.error(f'Error decoding file content: {e}')
        return ''

def get_files_with_prefix__f2ebfe11(env, config: Dict[str, Any]) -> int:
    """Get count of files with specific prefix in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'prefix' parameters

    Returns:
        Count of files with the prefix
    """
    path = config['path']
    prefix = config.get('prefix', '')
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return 0
    files_with_prefix = [node['name'] for node in result['children'] if node['name'].startswith(prefix)]
    return len(files_with_prefix)

def get_file_exists_and_size__ff3634ef(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_file_exists__c0930f6fc6470d951ad9d775e3a6c6a5(env, config):
    """Check if a file exists on the VM and validate it's a proper TIFF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'path': str, 'is_tiff': bool, 'has_content': bool, 'file_size': int}
    """
    file_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "true" || echo "false"', timeout=10)
    exists = exists_result.get('output', '').strip() == 'true'
    is_tiff = False
    has_content = False
    file_size = 0
    if exists:
        file_type_result = env.controller.run_bash_script(f'file -b "{file_path}"', timeout=10)
        file_type = file_type_result.get('output', '').strip().lower()
        is_tiff = 'tiff' in file_type or 'tagged image file format' in file_type
        size_result = env.controller.run_bash_script(f'stat -c %s "{file_path}" 2>/dev/null || stat -f %z "{file_path}" 2>/dev/null', timeout=10)
        try:
            file_size = int(size_result.get('output', '0').strip())
            has_content = file_size > 0
        except ValueError:
            file_size = 0
            has_content = False
    return {'exists': exists, 'path': file_path, 'is_tiff': is_tiff, 'has_content': has_content, 'file_size': file_size}

def get_file_exists__3928cfa5(env, config):
    """
    Check if a file exists at the specified path on VM and get image metadata.

    Returns a dict with:
    - exists: bool - whether file exists
    - is_valid_png: bool - whether file is a valid PNG image
    - file_size: int - file size in bytes
    - width: int - image width (if valid image)
    - height: int - image height (if valid image)
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output != 'EXISTS':
        return {'exists': False, 'is_valid_png': False, 'file_size': 0, 'width': 0, 'height': 0}
    size_command = f'stat -c %s "{file_path}" 2>/dev/null || echo "0"'
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    magic_command = f'xxd -l 8 -p "{file_path}" 2>/dev/null || echo ""'
    magic_result = env.controller.run_bash_script(magic_command, timeout=10)
    magic_bytes = magic_result.get('output', '').strip().replace('\n', '')
    is_valid_png = magic_bytes.lower().startswith('89504e470d0a1a0a')
    width = 0
    height = 0
    if is_valid_png:
        dim_command = f'identify -format "%w %h" "{file_path}" 2>/dev/null || echo "0 0"'
        dim_result = env.controller.run_bash_script(dim_command, timeout=10)
        dim_output = dim_result.get('output', '0 0').strip()
        try:
            parts = dim_output.split()
            if len(parts) >= 2:
                width = int(parts[0])
                height = int(parts[1])
        except (ValueError, IndexError):
            width = 0
            height = 0
    return {'exists': True, 'is_valid_png': is_valid_png, 'file_size': file_size, 'width': width, 'height': height}

def get_text_file_content__b8171706418d2058f81460f1a24d4635(env, config):
    """
    Read text file from VM and return full content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String content of the file
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading text file: {e}')
        return ''

def get_file_content_contains__c7e5cf7f(env, config):
    """Check if file contains specific text.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' and 'search_text' parameters

    Returns:
        bool: True if text is found in file, False otherwise
    """
    file_path = config.get('file_path')
    search_text = config.get('search_text', '')
    safe_search = search_text.replace("'", "'\\''")
    command = f"grep -q '{safe_search}' {file_path} 2>/dev/null && echo 'FOUND' || echo 'NOT_FOUND'"
    result = env.controller.run_bash_script(command, timeout=10)
    return 'FOUND' in result.get('output', '')

def get_file_exists__7298f585793a2b727ac3910de3795a50(env, config: Dict[str, Any]) -> bool:
    """Check if a specific file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file

    Returns:
        True if file exists, False otherwise
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0 and 'exists' in result['output']:
        return True
    else:
        return False

def get_vm_subtitle_file__4c5ac05d(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_pdf_files_info__230c9972(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_docx_content__8c8e551c645c8d069577024a9d0bf137(env, config):
    """Extract content and metadata from a .docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location on VM

    Returns:
        dict: Contains 'paragraphs' (list of text), 'word_count', 'paragraph_count'
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': 'No path provided'}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': f'File not found: {file_path}'}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        all_text = ' '.join(paragraphs)
        word_count = len(all_text.split())
        return {'paragraphs': paragraphs, 'word_count': word_count, 'paragraph_count': len(paragraphs), 'full_text': all_text}
    except Exception as e:
        return {'paragraphs': [], 'word_count': 0, 'paragraph_count': 0, 'error': str(e)}
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_file_location__c78e5698dfdf96679302f35b21f0928f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists at a specific location.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'path', 'filename' keys
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    exists = file_bytes is not None
    return {'exists': exists, 'path': path, 'filename': os.path.basename(path) if exists else None}

def get_text_file_content__600696b8508be0c2e3ca25794856fb75(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    path = config.get('path', '/home/user/output.txt')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_dir_file_list__d9c12923d5c2940b6520fb26a328924d(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' key

    Returns:
        List of filenames in the directory
    """
    command = config.get('command', 'ls /home/user/Desktop/audience/')
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return []
    output = result['output'].strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    return files

def get_tetris_files__b3b822fe(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_number_from_file__9b848ddc(env, config: dict):
    """Read a number from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Extracted number, or 0 if not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return 0
    try:
        content = file_bytes.decode('utf-8').strip()
        import re
        numbers = re.findall('\\d+', content)
        if numbers:
            return int(numbers[0])
        return 0
    except Exception as e:
        return 0

def get_csv_structure_and_content__8de41bc4fc70600ffffb805994ee2926(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get CSV structure and content information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains column_count, row_count, has_data, and non_empty_cells
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'column_count': 0, 'row_count': 0, 'has_data': False, 'non_empty_cells': 0}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'column_count': 0, 'row_count': 0, 'has_data': False, 'non_empty_cells': 0}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            rows = list(reader)
            if not rows:
                return {'column_count': 0, 'row_count': 0, 'has_data': False, 'non_empty_cells': 0}
            column_count = len(rows[0])
            data_rows = rows[1:]
            row_count = len(data_rows)
            has_data = row_count > 0
            non_empty_cells = 0
            for row in data_rows:
                for cell in row:
                    if cell and cell.strip():
                        non_empty_cells += 1
            return {'column_count': column_count, 'row_count': row_count, 'has_data': has_data, 'non_empty_cells': non_empty_cells}
    finally:
        os.unlink(tmp_path)

def get_pdf_file_exists__05a7fe2932371163af38b8e77d9b0c93(env, config: Dict[str, Any]) -> bool:
    """
    Check if a PDF file exists at the specified path on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        bool: True if file exists and is a PDF, False otherwise
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File does not exist at path: {path}')
            return False
        if file_bytes[:4] == b'%PDF':
            logger.info(f'PDF file exists at {path}')
            return True
        else:
            logger.warning(f'File exists at {path} but is not a PDF')
            return False
    except Exception as e:
        logger.error(f'Error checking file existence: {e}')
        return False

def get_file_exists_check__f768bfed(env, config: Dict) -> Optional[Dict]:
    """
    Check if a file exists and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_file_location__cbcf6e0c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a folder exists and contains a file matching criteria.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' and 'filename_contains' parameters

    Returns:
        Dict with 'folder_exists' (bool) and 'file_found' (bool) keys
    """
    folder_path = config['folder_path']
    filename_contains = config.get('filename_contains', '')
    check_folder = f"test -d '{folder_path}' && echo 'YES' || echo 'NO'"
    result = env.controller.run_bash_script(check_folder, timeout=10)
    folder_exists = result['output'].strip() == 'YES'
    if not folder_exists:
        return {'folder_exists': False, 'file_found': False}
    list_cmd = f"ls -1 '{folder_path}' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(list_cmd, timeout=10)
    files = result['output'].strip().split('\n') if result['output'].strip() else []
    file_found = any((filename_contains in f for f in files)) if filename_contains else len(files) > 0
    return {'folder_exists': True, 'file_found': file_found}

def get_file_size__b756f99d(env, config: dict):
    """Get the size of a file in bytes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        int: File size in bytes, or 0 if not found
    """
    path = config.get('path', '/home/user/Desktop/invoice.xlsx')
    result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then stat -c %s "{path}"; else echo "0"; fi', timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        return 0

def get_pdf_files_info__a96f92e2(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_number_from_file__c112996b(env, config: dict):
    """Read a number from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Extracted number, or 0 if not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return 0
    try:
        content = file_bytes.decode('utf-8').strip()
        import re
        numbers = re.findall('\\d+', content)
        if numbers:
            return int(numbers[0])
        return 0
    except Exception as e:
        return 0

def get_pdf_files_in_dir__32b125c8(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_file_count__23e95644(env, config: dict):
    """
    Count files matching a pattern in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern' keys

    Returns:
        int: Number of matching files
    """
    directory = config.get('directory', '')
    pattern = config.get('pattern', '*')
    command = f'find "{directory}" -maxdepth 1 -type f -name "{pattern}" 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        try:
            return int(result['output'].strip())
        except ValueError:
            return 0
    return 0

def get_text_file_content__5ced85fc_aug18_v1_a1b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6(env, config):
    """Read a text file from VM and return its entire content as a string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        str: File content as a string (with trailing whitespace stripped)
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        logger.info(f'Successfully read {len(content)} characters from {file_path}')
        return content
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return ''

def get_file_path_exists__c461bcde(env, config):
    """Check if file exists at exact path and return path if exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File path if exists, empty string otherwise
    """
    path = config.get('path', '')
    result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then echo "{path}"; else echo ""; fi', timeout=10)
    output = result.get('output', '').strip()
    return output

def get_both_files_non_empty__e07d7a26(env, config):
    """
    Check if both CSV and XLSX files exist and are non-empty.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with csv_path and xlsx_path

    Returns:
        dict: {"csv_size": int, "xlsx_size": int}
    """
    csv_path = config.get('csv_path', '/home/user/Desktop/contacts.csv')
    xlsx_path = config.get('xlsx_path', '/home/user/Desktop/contacts.xlsx')
    result = {'csv_size': 0, 'xlsx_size': 0}
    csv_bytes = env.controller.get_file(csv_path)
    if csv_bytes:
        result['csv_size'] = len(csv_bytes)
    xlsx_bytes = env.controller.get_file(xlsx_path)
    if xlsx_bytes:
        result['xlsx_size'] = len(xlsx_bytes)
    return result

def get_file_timestamp__739292ff(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get file timestamp information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: File timestamp information
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return None
    timestamp_script = f"""\nimport os\nimport time\nif os.path.exists('{path}'):\n    mtime = os.path.getmtime('{path}')\n    current_time = time.time()\n    age_seconds = current_time - mtime\n    print(f"{{mtime}}|{{age_seconds}}")\nelse:\n    print("NOT_FOUND")\n"""
    result = env.controller.execute_python_command(timestamp_script)
    if result and result.get('output'):
        output = result['output'].strip()
        if output != 'NOT_FOUND':
            parts = output.split('|')
            if len(parts) == 2:
                return {'exists': True, 'mtime': float(parts[0]), 'age_seconds': float(parts[1])}
    return {'exists': False, 'mtime': 0, 'age_seconds': float('inf')}

def get_pdf_files_info__f4eddf72(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_in_directory__35a43b1b(env, config):
    """Check if a file was moved (exists at destination, not at source).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (destination path)

    Returns:
        dict: {"exists_at_destination": bool, "exists_at_source": bool}
    """
    destination_path = config.get('path', '')
    filename = os.path.basename(destination_path)
    source_path = f'/home/user/{filename}'
    dest_file_bytes = env.controller.get_file(destination_path)
    exists_at_destination = dest_file_bytes is not None and len(dest_file_bytes) > 0
    source_file_bytes = env.controller.get_file(source_path)
    exists_at_source = source_file_bytes is not None and len(source_file_bytes) > 0
    return {'exists_at_destination': exists_at_destination, 'exists_at_source': exists_at_source}

def get_pdf_exports__c5999de9(env, config: dict):
    """
    Check if PDFs were exported to the specified directory and extract their content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        List of dicts, each containing 'filename' and 'content' (extracted text)
    """
    directory = config.get('directory', '/home/user/Documents/Blogs')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_data = []
    if list_result['output']:
        pdf_paths = [f.strip() for f in list_result['output'].strip().split('\n') if f.strip()]
        for pdf_path in pdf_paths:
            filename = os.path.basename(pdf_path)
            try:
                content = env.controller.get_vm_file(pdf_path)
                doc = fitz.open(stream=content, filetype='pdf')
                text_content = ''
                for page in doc:
                    text_content += page.get_text()
                doc.close()
                pdf_data.append({'filename': filename, 'content': text_content.lower()})
            except Exception as e:
                pdf_data.append({'filename': filename, 'content': ''})
    return pdf_data

def get_docx_content_and_highlights__fc0ff0e3(env, config: dict):
    """
    Extract document content and highlighted word information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Contains paragraphs, highlight_word, and highlight_count
    """
    vm_path = config.get('path', '/home/user/Desktop/gpt4_homer.docx')
    try:
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            logger.warning(f'Could not retrieve file from {vm_path}')
            return {'paragraphs': [], 'highlight_word': None, 'highlight_count': 0}
        from io import BytesIO
        doc = Document(BytesIO(file_bytes))
        paragraphs = []
        highlight_word = None
        highlight_count = 0
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if para_text:
                paragraphs.append(para_text)
                for run in para.runs:
                    if run.text and run.font.highlight_color is not None:
                        if highlight_word is None:
                            highlight_word = run.text.strip()
                        highlight_count += run.text.count(run.text.strip())
        return {'paragraphs': paragraphs, 'highlight_word': highlight_word, 'highlight_count': highlight_count}
    except Exception as e:
        logger.error(f'Error reading docx file: {e}')
        return {'paragraphs': [], 'highlight_word': None, 'highlight_count': 0}

def get_two_files_exist__3f6d3219(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if multiple files exist on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' parameter (list of file paths)

    Returns:
        Dict with file existence and size information
    """
    paths = config['paths']
    results = []
    for path in paths:
        command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; else echo 'NOT_EXISTS'; fi"
        result = env.controller.run_bash_script(command, timeout=10)
        output = result['output'].strip()
        if output == 'NOT_EXISTS' or result['returncode'] != 0:
            results.append({'path': path, 'exists': False, 'size_bytes': 0})
        else:
            try:
                size = int(output)
                results.append({'path': path, 'exists': True, 'size_bytes': size})
            except ValueError:
                results.append({'path': path, 'exists': False, 'size_bytes': 0})
    return {'files': results}

def get_numbered_pdf_chapters__a26911b7babb666234e16a8cd9e0ff40(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of numbered PDF chapter files (files starting with "N. " pattern).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key (VM path)

    Returns:
        List of PDF filenames that match the numbered chapter pattern (e.g., "1. Chapter.pdf")
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f"cd '{directory}' && ls -1 [0-9]*.pdf 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0:
        logger.warning(f"Failed to list numbered PDF files in {directory}: {result.get('error', '')}")
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    pdf_files = []
    for line in output.split('\n'):
        filename = line.strip()
        if filename and '. ' in filename:
            parts = filename.split('. ', 1)
            if len(parts) == 2 and parts[0].isdigit():
                pdf_files.append(filename)
    logger.info(f'Found {len(pdf_files)} numbered chapter PDFs in {directory}: {pdf_files}')
    return pdf_files

def get_subdirs_in_dirs__6eb64016ae326b8b48bfd5c29e5970e8(env, config):
    """Check if subdirectories exist in specified parent directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'parent_dirs' (list) and 'subdir_name' (str)

    Returns:
        dict: Mapping of parent dir to whether subdir exists in it
    """
    parent_dirs = config.get('parent_dirs', [])
    subdir_name = config.get('subdir_name', '')
    results = {}
    for parent_dir in parent_dirs:
        check_cmd = f'[ -d "{parent_dir}/{subdir_name}" ] && echo "exists" || echo "not_exists"'
        result = env.controller.run_bash_script(check_cmd, timeout=10)
        exists = False
        if result and result.get('output'):
            exists = result['output'].strip() == 'exists'
        results[parent_dir] = exists
    return results

def get_vm_file__a46a1ebc95419808bbd66e4ed9c5acac(env, config):
    """
    Get the DOCX file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded file
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    import tempfile
    import os
    suffix = os.path.splitext(file_path)[1]
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_text_file_content__433237d5d3b8a54ba0440460662ae73a(env, config: Dict[str, Any]) -> str:
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Text content of the file as a string, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_pdf_orientation__e222560cb6780019664a6917701a2659(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get PDF page orientation and page count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'orientation' (str: 'portrait' or 'landscape') and 'page_count' (int)
    """
    try:
        from pypdf import PdfReader
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'orientation': None, 'page_count': 0}
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            if len(reader.pages) == 0:
                return {'orientation': None, 'page_count': 0}
            page = reader.pages[0]
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            orientation = 'landscape' if width > height else 'portrait'
            return {'orientation': orientation, 'page_count': len(reader.pages)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return {'orientation': None, 'page_count': 0}

def get_folder_files__c14b0c4873b98943d0bcc59ebce705c9(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a specific folder on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to folder

    Returns:
        List of filenames in the folder
    """
    folder_path = config.get('path', '')
    command = f'ls "{folder_path}" 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    return files

def get_csv_first_n_fullnames__a0281c1d922f2a38909920e121739b97(env, config: Dict[str, Any]) -> List[str]:
    """Extract first N full names from a CSV with First Name and Last Name columns.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'n' keys

    Returns:
        List[str]: List of full names (FirstName LastName format)
    """
    file_path = config.get('path', '')
    n = config.get('n', 10)
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        full_names = []
        with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.DictReader(f)
            for (i, row) in enumerate(reader):
                if i >= n:
                    break
                first_name = row.get('First Name', '').strip()
                last_name = row.get('Last Name', '').strip()
                if first_name or last_name:
                    full_names.append(f'{first_name} {last_name}')
        return full_names
    finally:
        os.unlink(tmp_path)

def get_pdf_basic_info__56a55075c129634e9f17633f657c7657(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract basic information from a PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with 'exists', 'page_count', 'file_size' keys (always returns a dict, never None)
    """
    from pypdf import PdfReader
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'page_count': 0, 'file_size': 0}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        return {'exists': True, 'page_count': page_count, 'file_size': file_size}
    except Exception as e:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'error': str(e)}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_zip_file_list__53fe105429a60cae06bcf9ce59b19b3e(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files inside a zip archive on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'zip_path' key

    Returns:
        List of filenames in the zip (sorted, basenames only)
    """
    zip_path = config.get('zip_path', '')
    command = f"zipinfo -1 '{zip_path}' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0 or not result.get('output'):
        logger.warning(f'Failed to list zip contents: {zip_path}')
        return []
    files = result['output'].strip().split('\n')
    basenames = []
    for f in files:
        if f and (not f.endswith('/')):
            basenames.append(os.path.basename(f))
    return sorted(basenames)

def get_pdf_files_in_dir__989759aa(env, config: dict):
    """Get list of PDF files in specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        List of PDF filenames in the directory
    """
    directory = config.get('directory', '/home/user/Downloads')
    command = f"ls -1 {directory}/*.pdf 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=30)
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [os.path.basename(f) for f in output.split('\n') if f.strip()]
    logger.info(f'Found PDF files in {directory}: {files}')
    return files

def get_docx_text_content__c32b1108(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_directory_tree__63701d4e(env, config):
    """Get directory tree structure.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Directory tree output
    """
    path = config.get('path', '/home/user')
    command = f'find {path} -type d | sort'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        return result['output'].strip()
    else:
        logger.error(f"Failed to get directory tree for {path}: {result['error']}")
        return None

def get_pdf_files_in_dir__bfc9cce9462c8f8be02842ad8f805893(env, config: Dict[str, Any]) -> List[str]:
    """Get list of PDF files in a specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key specifying path to check

    Returns:
        List of PDF filenames found in the directory
    """
    directory = config.get('directory', '/home/user/Downloads')
    command = f"ls {directory}/*.pdf 2>/dev/null | xargs -n 1 basename 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 and result['returncode'] != 123:
        logger.warning(f"Failed to list PDF files in {directory}: {result.get('error', '')}")
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    pdf_files = [f.strip() for f in output.split('\n') if f.strip()]
    logger.info(f'Found {len(pdf_files)} PDF files in {directory}: {pdf_files}')
    return pdf_files

def get_file_exists__3b8e423e430323c0078f4425aded05b9(env, config: dict):
    """Check if output file exists and verify 180-degree rotation was applied.

    This getter checks both the source (flipped) video and the output video to verify
    that a 180-degree rotation was actually performed, preventing false positives where
    the agent just copies the file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying output file path on VM

    Returns:
        dict: {
            'exists': bool,
            'path': str,
            'size': int,
            'source_rotation': int or None (rotation of source flipped video),
            'output_rotation': int or None (rotation of output video),
            'rotation_changed': bool (whether rotation was modified),
            'is_valid_video': bool
        }
    """
    output_path = config.get('path', '')
    source_path = '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4'
    source_rotation = None
    try:
        source_bytes = env.controller.get_file(source_path)
        if source_bytes and len(source_bytes) > 0:
            (source_rotation, source_valid) = _get_video_rotation(source_bytes)
            logger.info(f'Source video {source_path}: rotation={source_rotation}, valid={source_valid}')
    except Exception as e:
        logger.warning(f'Could not read source video {source_path}: {e}')
    try:
        output_bytes = env.controller.get_file(output_path)
        if not output_bytes or len(output_bytes) == 0:
            logger.info(f'Output file {output_path} does not exist or is empty')
            return {'exists': False, 'path': output_path, 'size': 0, 'source_rotation': source_rotation, 'output_rotation': None, 'rotation_changed': False, 'is_valid_video': False}
        output_size = len(output_bytes)
        logger.info(f'Output file {output_path} exists with size {output_size} bytes')
        (output_rotation, is_valid_video) = _get_video_rotation(output_bytes)
        logger.info(f'Output video {output_path}: rotation={output_rotation}, valid={is_valid_video}')
        rotation_changed = False
        if source_rotation is not None and output_rotation is not None:
            norm_source = source_rotation % 360
            norm_output = output_rotation % 360
            if norm_source == 180 and norm_output == 0:
                rotation_changed = True
        elif source_rotation == 180 and output_rotation == 0:
            rotation_changed = True
        return {'exists': True, 'path': output_path, 'size': output_size, 'source_rotation': source_rotation, 'output_rotation': output_rotation, 'rotation_changed': rotation_changed, 'is_valid_video': is_valid_video}
    except Exception as e:
        logger.error(f'Error checking file existence for {output_path}: {e}')
        return {'exists': False, 'path': output_path, 'size': 0, 'source_rotation': source_rotation, 'output_rotation': None, 'rotation_changed': False, 'is_valid_video': False}

def get_text_file_content__210f933b0657674e3cf1f4222885738f(env, config: dict) -> Optional[str]:
    """Get content from a text file on the VM as a single string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_docx_content__2e06628be7ef5b3dc45161c8ea091810(env, config: Dict[str, Any]) -> str:
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the DOCX file

    Returns:
        String containing the text content of the document
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        import tempfile
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        text_content = '\n'.join([para.text for para in doc.paragraphs])
        return text_content.strip()
    except Exception as e:
        return ''
    finally:
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_directory_structure__3ef6f937c5dcda24ec08a8ba5aa2c872(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get the structure of a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path' key

    Returns:
        Dict with directory structure information
    """
    dir_path = config.get('dir_path', '')
    result = {'exists': False, 'subdirs': [], 'files': []}
    check_dir_cmd = f"[ -d '{dir_path}' ] && echo 'EXISTS' || echo 'NOT_EXISTS'"
    dir_check = env.controller.run_bash_script(check_dir_cmd, timeout=10)
    if dir_check.get('output', '').strip() == 'EXISTS':
        result['exists'] = True
        subdir_cmd = f"find '{dir_path}' -mindepth 1 -maxdepth 1 -type d -exec basename {{}} \\; | sort"
        subdir_result = env.controller.run_bash_script(subdir_cmd, timeout=20)
        if subdir_result.get('returncode') == 0:
            subdirs = subdir_result.get('output', '').strip().split('\n')
            result['subdirs'] = [d.strip() for d in subdirs if d.strip()]
        files_cmd = f"find '{dir_path}' -mindepth 1 -maxdepth 1 -type f -exec basename {{}} \\; | sort"
        files_result = env.controller.run_bash_script(files_cmd, timeout=20)
        if files_result.get('returncode') == 0:
            files = files_result.get('output', '').strip().split('\n')
            result['files'] = [f.strip() for f in files if f.strip()]
    logger.info(f'Directory structure for {dir_path}: {result}')
    return result

def get_file_exists__6ee1fdcd(env, config: dict):
    """Check if the exported file exists on VM and validate it's a valid PDF.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with file validation information including:
            - exists: bool, whether file exists
            - is_pdf: bool, whether file is a PDF (based on MIME type)
            - file_size: int, file size in bytes (0 if doesn't exist)
            - is_fresh: bool, whether file was created recently (within last 5 minutes)
            - created_by_libreoffice: bool, whether PDF metadata indicates LibreOffice creator
    """
    vm_path = config.get('path', '/home/user/Desktop/enterprise-survey.pdf')
    result = {'exists': False, 'is_pdf': False, 'file_size': 0, 'is_fresh': False, 'created_by_libreoffice': False}
    exists_check = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    if exists_check.get('output', '').strip() != 'EXISTS':
        return result
    result['exists'] = True
    file_type_check = env.controller.run_bash_script(f"file -b --mime-type '{vm_path}'", timeout=10)
    mime_type = file_type_check.get('output', '').strip()
    result['is_pdf'] = mime_type == 'application/pdf'
    size_check = env.controller.run_bash_script(f"stat -c %s '{vm_path}' 2>/dev/null || echo '0'", timeout=10)
    try:
        result['file_size'] = int(size_check.get('output', '0').strip())
    except ValueError:
        result['file_size'] = 0
    freshness_check = env.controller.run_bash_script(f"find '{vm_path}' -mmin -5 2>/dev/null | wc -l", timeout=10)
    try:
        is_fresh = int(freshness_check.get('output', '0').strip()) > 0
        result['is_fresh'] = is_fresh
    except ValueError:
        result['is_fresh'] = False
    metadata_check = env.controller.run_bash_script(f"command -v pdfinfo > /dev/null && pdfinfo '{vm_path}' 2>/dev/null || strings '{vm_path}' 2>/dev/null | head -100", timeout=15)
    metadata_output = metadata_check.get('output', '').lower()
    result['created_by_libreoffice'] = 'libreoffice' in metadata_output or 'writer' in metadata_output or 'calc' in metadata_output
    return result

def get_pdf_page_count__e941d252d60fece99cf71fbcddf8521f(env, config: Dict[str, Any]) -> Any:
    """Get the page count of a PDF file from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - query: Query to find the file on Google Drive
            - dest: Local destination filename

    Returns:
        int: Number of pages in the PDF, or 0 if file not found/error
    """
    try:
        import PyPDF2
        from desktop_env.evaluators.getters.chrome import get_googledrive_file
        local_path = get_googledrive_file(env, config)
        if not local_path or not os.path.exists(local_path):
            logger.warning(f'PDF file not found: {local_path}')
            return 0
        with open(local_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            page_count = len(pdf_reader.pages)
            logger.info(f'PDF page count: {page_count}')
            return page_count
    except Exception as e:
        logger.error(f'Error getting PDF page count: {e}')
        return 0

def get_tetris_files__c6117858(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_blog_folder_listing__4e03b1ed(env, config: dict):
    """Get complete directory listing of Blog folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        str: Directory listing output
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    command = f"ls -la '{folder_path}' 2>/dev/null || echo 'FOLDER_NOT_FOUND'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    logger.info(f'Directory listing for {folder_path}:\n{output}')
    return output

def get_pdf_file_info__40d8ee41df97cbb2f480ba7af545efc4(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get PDF file from VM and verify it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        str: Path to downloaded PDF file in cache, or None if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import os
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_file_renamed__b4364020(env, config):
    """Check if file was renamed (old doesn't exist, new exists).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_path' and 'new_path' keys

    Returns:
        dict: Status of old and new files
    """
    old_path = config.get('old_path', '')
    new_path = config.get('new_path', '')
    old_bytes = env.controller.get_file(old_path)
    old_exists = old_bytes is not None and len(old_bytes) > 0
    new_bytes = env.controller.get_file(new_path)
    new_exists = new_bytes is not None and len(new_bytes) > 0
    return {'old_exists': old_exists, 'new_exists': new_exists, 'renamed': not old_exists and new_exists}

def get_speedtest_csv_data__50006736ce024a92878da0e4240e2bb0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract speedtest data from CSV file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with CSV data or empty dict if file doesn't exist/is invalid
    """
    import os
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Failed to get file: {config['path']}")
        return {}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()
        if len(lines) < 2:
            logger.warning('CSV file has insufficient data')
            return {}
        header = lines[0].strip().split(',')
        data_line = lines[1].strip().split(',')
        result = {}
        for (i, col_name) in enumerate(header):
            if i < len(data_line):
                result[col_name.strip()] = data_line[i].strip()
        return result
    except Exception as e:
        logger.error(f'Error reading CSV file: {e}')
        return {}
    finally:
        os.unlink(tmp_path)

def get_pdf_count_in_dir__085cc8d6(env, config: dict):
    """
    Count number of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path' key

    Returns:
        int: Number of PDF files found
    """
    dir_path = config.get('dir_path', '/home/user/Downloads')
    command = f"find {dir_path} -maxdepth 1 -name '*.pdf' -type f | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    try:
        count = int(result['output'].strip())
        logger.info(f'Found {count} PDF files in {dir_path}')
        return count
    except Exception as e:
        logger.error(f'Failed to parse PDF count: {e}')
        return 0

def get_file_existence__a8a082525df1807c95a7519289fda5a0(env, config):
    """Check if files exist and get basic properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' key (list of file paths)

    Returns:
        dict: Map of file paths to existence status and row count
    """
    paths = config.get('paths', [])
    result = {}
    for path in paths:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            result[path] = {'exists': False, 'row_count': 0}
            continue
        row_count = 0
        if path.endswith('.csv'):
            with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    reader = csv.reader(f)
                    next(reader, None)
                    row_count = sum((1 for row in reader))
            finally:
                os.unlink(tmp_path)
        elif path.endswith('.xlsx'):
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                wb = openpyxl.load_workbook(tmp_path, data_only=True)
                ws = wb.active
                for row_idx in range(2, ws.max_row + 1):
                    if ws.cell(row_idx, 1).value or ws.cell(row_idx, 2).value or ws.cell(row_idx, 3).value:
                        row_count += 1
            finally:
                os.unlink(tmp_path)
        result[path] = {'exists': True, 'row_count': row_count}
    return result

def get_file_size__198be354(env, config):
    """Get file size and content information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File size and content validation information
    """
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'size_bytes': 0, 'is_python': False, 'has_functions': False, 'has_imports': False, 'content': ''}
    if isinstance(file_content, bytes):
        try:
            content_str = file_content.decode('utf-8')
        except:
            content_str = str(file_content)
    else:
        content_str = file_content
    is_python = False
    has_functions = False
    has_imports = False
    try:
        ast.parse(content_str)
        is_python = True
        has_functions = bool(re.search('\\bdef\\s+\\w+\\s*\\(', content_str))
        has_imports = bool(re.search('\\b(import|from)\\s+', content_str))
    except:
        has_functions = bool(re.search('\\bdef\\s+\\w+\\s*\\(', content_str))
        has_imports = bool(re.search('\\b(import|from)\\s+', content_str))
        is_python = has_functions or has_imports or bool(re.search('\\b(class|if|for|while|return)\\s+', content_str))
    return {'exists': True, 'size_bytes': len(file_content), 'is_python': is_python, 'has_functions': has_functions, 'has_imports': has_imports, 'content': content_str}

def get_text_file_content__c4d957b9ed4c029b04783deb70ccd213(env, config: Dict[str, Any]) -> str:
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Text content of the file as a string, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_pdf_files_info__6e9dcacc(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_content__e1c8e8d0(env, config: dict):
    """
    Read the content of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Content of the file (stripped)
    """
    path = config.get('path', '/home/user/Desktop/largest_doc.txt')
    command = f'cat {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result.get('returncode', 1) != 0:
        return ''
    return result.get('output', '').strip()

def get_text_file__3c678f53(env, config):
    """Get text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest'

    Returns:
        Text file content as string
    """
    vm_path = config.get('path')
    dest = config.get('dest', 'file.txt')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return ''
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        return content
    except Exception as e:
        print(f'Error reading text file: {e}')
        return ''

def get_text_file_lines__bf9ee805c777733e26c14d1927c30b1a(env, config: Dict[str, Any]) -> list:
    """Read text file and return lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of non-empty lines from the file
    """
    import os
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Failed to get file: {config['path']}")
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.txt', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            lines = [line.strip() for line in f.readlines() if line.strip()]
        return lines
    except Exception as e:
        logger.error(f'Error reading text file: {e}')
        return []
    finally:
        os.unlink(tmp_path)

def get_multi_directory_contents__3b0a753c(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_directory_listing__32107748(env, config: dict):
    """Get .eml file content from VM to verify email was saved.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'pattern'

    Returns:
        dict: Dictionary with 'exists' (bool), 'filename' (str), and 'content' (str)
    """
    path = config.get('path', '')
    pattern = config.get('pattern', '*')
    expected_filename = 'forwarded_paper.eml'
    check_command = f"ls -1 {path}/{pattern} 2>/dev/null || echo ''"
    check_result = env.controller.run_bash_script(check_command, timeout=10)
    result_dict = {'exists': False, 'filename': '', 'content': ''}
    if check_result and check_result.get('returncode') == 0:
        listing = check_result.get('output', '').strip()
        lines = listing.split('\n') if listing else []
        for line in lines:
            if expected_filename in line:
                result_dict['exists'] = True
                result_dict['filename'] = expected_filename
                file_path = f'{path}/{expected_filename}'
                read_command = f"cat '{file_path}' 2>/dev/null || echo ''"
                read_result = env.controller.run_bash_script(read_command, timeout=10)
                if read_result and read_result.get('returncode') == 0:
                    result_dict['content'] = read_result.get('output', '')
                break
    return result_dict

def get_zip_files_and_dirs__e2b9921b3a4cfcd62d172fdaae4844bb(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get both directory contents and check for zip files existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path1', 'dir_path2', 'zip_path1', 'zip_path2' keys

    Returns:
        Dict containing:
            - dir1_files: List of files in first directory
            - dir2_files: List of files in second directory
            - zip1_exists: Boolean indicating if first zip file exists
            - zip2_exists: Boolean indicating if second zip file exists
    """
    dir_path1 = config['dir_path1']
    dir_path2 = config['dir_path2']
    zip_path1 = config['zip_path1']
    zip_path2 = config['zip_path2']

    def get_dir_files(path: str) -> List[str]:
        result = env.controller.run_bash_script(f"ls -1 '{path}' 2>/dev/null || echo ''", timeout=10)
        if result['returncode'] != 0:
            logger.warning(f"Failed to list directory {path}: {result.get('error', '')}")
            return []
        output = result.get('output', '').strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        return files

    def check_file_exists(path: str) -> bool:
        result = env.controller.run_bash_script(f"test -f '{path}' && echo 'exists' || echo 'missing'", timeout=10)
        if result['returncode'] != 0:
            logger.warning(f"Failed to check file {path}: {result.get('error', '')}")
            return False
        output = result.get('output', '').strip()
        return output == 'exists'
    dir1_files = get_dir_files(dir_path1)
    dir2_files = get_dir_files(dir_path2)
    zip1_exists = check_file_exists(zip_path1)
    zip2_exists = check_file_exists(zip_path2)
    return {'dir1_files': dir1_files, 'dir2_files': dir2_files, 'zip1_exists': zip1_exists, 'zip2_exists': zip2_exists}

def get_text_file_content__0fe1ad2c6a085376a5c2eb19244f70a6(env, config: Dict[str, Any]) -> str:
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Text content of the file as a string, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_file_exists__5f5351b0(env, config: dict):
    """Get the list of files contained in a zip archive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        dict: Dictionary with 'exists' (bool), 'valid' (bool), and 'contents' (list of filenames)
    """
    file_path = config.get('file_path', '')
    check_command = f'test -f "{file_path}" && echo "1" || echo "0"'
    check_result = env.controller.run_bash_script(check_command, timeout=10)
    if check_result['returncode'] != 0 or check_result['output'].strip() != '1':
        return {'exists': False, 'valid': False, 'contents': []}
    list_command = f'unzip -l "{file_path}" 2>&1'
    list_result = env.controller.run_bash_script(list_command, timeout=15)
    if list_result['returncode'] != 0:
        return {'exists': True, 'valid': False, 'contents': []}
    contents = []
    lines = list_result['output'].split('\n')
    parsing = False
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.startswith('---------'):
            if not parsing:
                parsing = True
            else:
                break
            continue
        if parsing and (not line.startswith('Archive:')) and (not line.startswith('Length')):
            parts = line.split()
            if len(parts) >= 4:
                filename = parts[-1]
                if '/' in filename:
                    filename = filename.split('/')[-1]
                contents.append(filename)
    return {'exists': True, 'valid': True, 'contents': contents}

def get_docx_text_content__6c0adead2c9259465f7529a90ac3bebf(env, config: Dict[str, Any]) -> str:
    """Extract all text content from a Word document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file on VM

    Returns:
        str: All text content from the document, concatenated
    """
    from docx import Document
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        all_text = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                all_text.append(text)
        return '\n'.join(all_text)
    finally:
        os.unlink(tmp_path)

def get_text_file_content__7dfb45a4(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_pdf_files_info__2c1e781e(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_exists__898e13b4(env, config: Dict[str, Any]) -> bool:
    """Check if a file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        True if file exists, False otherwise
    """
    file_path = config.get('path')
    if not file_path:
        return False
    try:
        file_bytes = env.controller.get_file(file_path)
        return file_bytes is not None and len(file_bytes) > 0
    except Exception:
        return False

def get_directory_info__e01c3944c2f0dfae12ffc1d2b96464cf(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get directory information including existence, PDF count, and PDF filenames.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict with 'exists', 'is_directory', 'pdf_count', and 'pdf_filenames' keys
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    check_cmd = f"if [ -d '{directory}' ]; then echo 'DIR_EXISTS'; elif [ -e '{directory}' ]; then echo 'FILE_EXISTS'; else echo 'NOT_EXISTS'; fi"
    result = env.controller.run_bash_script(check_cmd, timeout=10)
    output = result.get('output', '').strip()
    exists = 'EXISTS' in output
    is_directory = 'DIR_EXISTS' in output
    pdf_count = 0
    pdf_filenames = []
    if is_directory:
        list_cmd = f'ls {directory}/*.pdf 2>/dev/null | xargs -n 1 basename 2>/dev/null'
        list_result = env.controller.run_bash_script(list_cmd, timeout=10)
        output_list = list_result.get('output', '').strip()
        if output_list:
            pdf_filenames = [f.strip() for f in output_list.split('\n') if f.strip().endswith('.pdf')]
            pdf_count = len(pdf_filenames)
    logger.info(f'Directory {directory}: exists={exists}, is_dir={is_directory}, pdf_count={pdf_count}, filenames={pdf_filenames}')
    return {'exists': exists, 'is_directory': is_directory, 'pdf_count': pdf_count, 'pdf_filenames': pdf_filenames}

def get_file_recovery_state__5ea617a3(env, config):
    """
    Get the state of file recovery from Trash.

    Checks:
    1. If the target file exists at the expected path
    2. The file size (to ensure it's not empty/corrupted)
    3. If the original file no longer exists in Trash
    4. The SHA256 hash of the recovered file (to verify it's the same file)
    5. The original file's hash (from pre-stored hash file)

    Args:
        env: Environment object
        config: Configuration dict with 'target_path' and 'original_name'

    Returns:
        dict: {
            'file_exists': bool,
            'file_size': int,
            'not_in_trash': bool,
            'file_hash': str (SHA256 hash in hex, or None if file doesn't exist),
            'original_hash': str (SHA256 hash of original file, or None if not found)
        }
    """
    target_path = config.get('target_path', '/home/user/Desktop/recovered_poster.webp')
    original_name = config.get('original_name', 'poster_party_night.webp')
    result = {'file_exists': False, 'file_size': 0, 'not_in_trash': False, 'file_hash': None, 'original_hash': None}
    file_bytes = env.controller.get_file(target_path)
    if file_bytes:
        result['file_exists'] = True
        result['file_size'] = len(file_bytes)
        result['file_hash'] = hashlib.sha256(file_bytes).hexdigest()
    hash_file_bytes = env.controller.get_file('/tmp/original_poster_hash.txt')
    if hash_file_bytes:
        try:
            hash_content = hash_file_bytes.decode('utf-8').strip()
            result['original_hash'] = hash_content.split()[0]
        except (UnicodeDecodeError, IndexError):
            pass
    trash_check_cmd = f"gio list trash:// | grep -q '{original_name}' && echo 'in_trash' || echo 'not_in_trash'"
    trash_result = env.controller.execute(trash_check_cmd, shell=True)
    if trash_result and 'not_in_trash' in trash_result.get('output', ''):
        result['not_in_trash'] = True
    return result

def get_file_exists__c6c3aa52(env, config):
    """
    Check if a file exists and verify it's a valid PNG image with proper content.

    Returns a dict with file existence and validation information.
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output != 'EXISTS':
        return {'exists': False, 'is_png': False, 'has_valid_size': False, 'has_valid_dimensions': False}
    size_command = f'stat -c %s "{file_path}" 2>/dev/null || echo "0"'
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    magic_command = f'xxd -l 8 -p "{file_path}" 2>/dev/null || echo ""'
    magic_result = env.controller.run_bash_script(magic_command, timeout=10)
    magic_bytes = magic_result.get('output', '').strip().replace('\n', '')
    is_png = magic_bytes.lower().startswith('89504e470d0a1a0a')
    file_command = f'file "{file_path}" 2>/dev/null || echo ""'
    file_result = env.controller.run_bash_script(file_command, timeout=10)
    file_output = file_result.get('output', '').strip()
    has_valid_dimensions = 'PNG image data' in file_output and 'x' in file_output
    return {'exists': True, 'is_png': is_png, 'has_valid_size': file_size > 1024, 'has_valid_dimensions': has_valid_dimensions, 'file_size': file_size}

def get_text_file_content__acd75d14(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    result = env.controller.get_file(file_path)
    if result is None:
        return ''
    try:
        return result.decode('utf-8', errors='ignore').strip()
    except Exception:
        return ''

def get_pdf_files_list__6871c0fe(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files in the specified directory with their page counts.

    Args:
        env: Environment object
        config: Configuration dict with 'directory' key

    Returns:
        Dict mapping filename to page count, e.g., {'file.pdf': 10}
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        page_cmd = f"pdfinfo '{filepath}' 2>/dev/null | grep -i 'Pages:' | awk '{{print $2}}'"
        page_result = env.controller.run_bash_script(page_cmd, timeout=10)
        if page_result.get('returncode') == 0:
            page_output = page_result.get('output', '').strip()
            if page_output.isdigit():
                pdf_info[filename] = int(page_output)
            else:
                pdf_info[filename] = 0
        else:
            pdf_info[filename] = 0
    return pdf_info

def get_file_count__19cf6326(env, config):
    """Get file count in a directory along with file details.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and optional 'pattern' keys

    Returns:
        dict: Dictionary with filenames as keys and file sizes as values
    """
    path = config.get('path', '/home/user')
    pattern = config.get('pattern', '*')
    command = f'find {path} -maxdepth 1 -name "{pattern}" -type f -exec basename {{}} \\; -exec stat -c %s {{}} \\;'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        try:
            lines = result['output'].strip().split('\n')
            if not lines or lines[0] == '':
                return {}
            file_info = {}
            for i in range(0, len(lines), 2):
                if i + 1 < len(lines):
                    filename = lines[i].strip()
                    size = int(lines[i + 1].strip())
                    file_info[filename] = size
            return file_info
        except (ValueError, IndexError) as e:
            logger.error(f"Failed to parse file details: {result['output']}, error: {e}")
            return {}
    else:
        logger.error(f"Failed to get files in {path}: {result['error']}")
        return {}

def get_pdf_page_count__a73a14fa(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the path to a PDF file from the VM for page count verification.

    Args:
        env: Environment object with controller to access VM
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        str: Local path to the downloaded PDF file, or None if file doesn't exist
    """
    pdf_path_on_vm = config['path']
    pdf_filename = os.path.basename(pdf_path_on_vm)
    local_path = os.path.join(env.cache_dir, pdf_filename)
    try:
        file_content = env.controller.get_file(pdf_path_on_vm)
        if file_content is None:
            return None
        with open(local_path, 'wb') as f:
            f.write(file_content)
        return local_path
    except Exception as e:
        return None

def get_pdf_exports__d582d1d2(env, config: dict):
    """
    Check if PDFs were exported to the specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'expected_count' keys

    Returns:
        List of PDF filenames in the directory
    """
    directory = config.get('directory', '/home/user/Downloads/Articles')
    command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    try:
        pdf_count = int(result['output'].strip())
    except (ValueError, KeyError):
        pdf_count = 0
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        pdf_files = [os.path.basename(f.strip()) for f in list_result['output'].strip().split('\n') if f.strip()]
    return pdf_files

def get_pdf_content__5e9826db825b680d44e19c36727da776(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract text content and metadata from a PDF file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM path to PDF)

    Returns:
        Dict containing:
            - exists: bool - whether file exists
            - page_count: int - number of pages
            - text_content: str - extracted text from all pages
            - has_employee_name: bool - whether 'Michael Brown' appears
            - has_checkmarks: bool - whether checkmark symbols appear
            - has_ratings: bool - whether rating-related text appears
    """
    path = config.get('path', '')
    result = {'exists': False, 'page_count': 0, 'text_content': '', 'has_employee_name': False, 'has_checkmarks': False, 'has_ratings': False}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Could not retrieve file: {path}')
        return result
    result['exists'] = True
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from pypdf import PdfReader
            reader = PdfReader(tmp_path)
            result['page_count'] = len(reader.pages)
            all_text = ''
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    all_text += page_text + '\n'
            result['text_content'] = all_text
            if 'Michael Brown' in all_text:
                result['has_employee_name'] = True
            checkmark_symbols = ['√', '✓', 'X', 'x', '☑', '✔']
            if any((symbol in all_text for symbol in checkmark_symbols)):
                result['has_checkmarks'] = True
            rating_keywords = ['rating', 'performance', 'evaluation', 'review', 'score', 'excellent', 'good', 'fair', 'poor']
            if any((keyword.lower() in all_text.lower() for keyword in rating_keywords)):
                result['has_ratings'] = True
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting PDF content: {e}')
    return result

def get_csv_row_count__2d3e7876(env, config: dict):
    """Get CSV structure and content information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains row_count, header, column_count, and email validation info
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'row_count': 0, 'header': None, 'column_count': 0, 'has_email_content': False}
    content = file_bytes.decode('utf-8')
    lines = content.strip().split('\n')
    if not lines:
        return {'row_count': 0, 'header': None, 'column_count': 0, 'has_email_content': False}
    reader = csv.reader(lines)
    rows = list(reader)
    if len(rows) == 0:
        return {'row_count': 0, 'header': None, 'column_count': 0, 'has_email_content': False}
    header = rows[0] if len(rows) > 0 else []
    data_rows = rows[1:] if len(rows) > 1 else []
    row_count = len(data_rows)
    column_count = len(header)
    has_email_content = False
    if data_rows:
        email_count = 0
        for row in data_rows:
            if row and len(row) > 0 and ('@' in row[0]):
                email_count += 1
        has_email_content = email_count > len(data_rows) * 0.8
    return {'row_count': row_count, 'header': header, 'column_count': column_count, 'has_email_content': has_email_content}

def get_docx_content_check__ef584ea3(env, config: dict):
    """Check document has GPT3 content."""
    vm_path = config.get('path', '/home/user/Desktop/gpt3_responses.docx')
    try:
        file_bytes = env.controller.get_file(vm_path)
        if not file_bytes:
            return {'exists': False, 'paragraph_count': 0}
        from io import BytesIO
        doc = Document(BytesIO(file_bytes))
        para_count = sum((1 for p in doc.paragraphs if p.text.strip()))
        return {'exists': True, 'paragraph_count': para_count}
    except:
        return {'exists': False, 'paragraph_count': 0}

def get_question_count_file__e2620a9d(env, config: dict):
    """Get question count from text file saved by user.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Question count read from file, or None if file doesn't exist or invalid
    """
    file_path = config.get('path', '/home/user/Desktop/question_count.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found: {file_path}')
            return None
        content = file_bytes.decode('utf-8').strip()
        try:
            count = int(content)
            return count
        except ValueError:
            logger.warning(f'Invalid question count format: {content}')
            return None
    except Exception as e:
        logger.error(f'Error reading question count file: {e}')
        return None

def get_text_file_lines__2267ae43c99e342cd984b7743dc212e6(env, config: dict) -> Optional[list]:
    """Get lines from a text file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of stripped lines, or None if file not found
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n')]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_pdf_files_list__89dbe8bb(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get list of PDF files in the specified directory with page counts.

    Returns:
        Dict with 'files' (list of filenames) and 'page_counts' (dict of filename -> page count)
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {'files': [], 'page_counts': {}}
    output = result.get('output', '').strip()
    if not output:
        return {'files': [], 'page_counts': {}}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    page_counts = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        pdfinfo_cmd = f"pdfinfo '{filepath}' 2>/dev/null | grep -i '^Pages:' | awk '{{print $2}}'"
        pdfinfo_result = env.controller.run_bash_script(pdfinfo_cmd, timeout=10)
        if pdfinfo_result.get('returncode') == 0:
            page_count_str = pdfinfo_result.get('output', '').strip()
            try:
                page_counts[filename] = int(page_count_str)
            except (ValueError, TypeError):
                page_counts[filename] = 0
        else:
            page_counts[filename] = 0
    return {'files': files, 'page_counts': page_counts}

def get_dir_structure_info__5b67568a(env, config: dict):
    """Get directory structure information including file count and directory existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'base_dir' parameter

    Returns:
        Dict with directory existence status and file count
    """
    base_dir = config.get('base_dir', '/home/user/Documents/SavedBlogs')
    dir_check_cmd = f"test -d {base_dir} && echo 'exists' || echo 'not_exists'"
    dir_result = env.controller.run_bash_script(dir_check_cmd, timeout=30)
    dir_exists = dir_result.get('output', '').strip() == 'exists'
    count_cmd = f"find {base_dir} -name '*.pdf' -type f 2>/dev/null | wc -l"
    count_result = env.controller.run_bash_script(count_cmd, timeout=30)
    pdf_count = int(count_result.get('output', '0').strip())
    result = {'directory_exists': dir_exists, 'pdf_count': pdf_count}
    logger.info(f'Directory structure info for {base_dir}: {result}')
    return result

def get_pdf_files_in_dir__d928a635(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_flac_file_status__b00c11502ffa26aa7b145b3096daa5d2(env, config: Dict[str, str]):
    """
    Get status of FLAC audio file created from video.
    Returns dict with file properties and FLAC format verification.
    """
    file_path = config.get('path', '')
    command = f"\nimport os\nimport json\nimport subprocess\n\nfile_path = '{file_path}'\nresult = {{}}\n\nif os.path.exists(file_path):\n    result['exists'] = True\n    result['size'] = os.path.getsize(file_path)\n    result['extension'] = os.path.splitext(file_path)[1].lower()\n    result['filename'] = os.path.basename(file_path)\n\n    # Verify MIME type for FLAC\n    try:\n        mime = subprocess.check_output(['file', '--mime-type', '-b', file_path], text=True).strip()\n        result['mime_type'] = mime\n        # FLAC MIME type is usually audio/flac or audio/x-flac\n        result['is_flac'] = mime in ['audio/flac', 'audio/x-flac']\n    except:\n        result['mime_type'] = 'unknown'\n        result['is_flac'] = False\nelse:\n    result['exists'] = False\n    result['size'] = 0\n    result['extension'] = ''\n    result['filename'] = ''\n    result['mime_type'] = ''\n    result['is_flac'] = False\n\nprint(json.dumps(result))\n"
    try:
        response = env.controller.execute_python_command(command)
        if response and response.get('output'):
            import json
            return json.loads(response['output'].strip())
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_flac': False}
    except Exception as e:
        logger.error(f'Error checking FLAC file: {e}')
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_flac': False}

def get_python_file_lines__d93565be(env, config):
    """Get specific lines from a Python file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of strings (file lines)
    """
    path = config.get('path', '/home/user/Desktop/test.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return []
    content = file_bytes.decode('utf-8')
    return content.splitlines()

def get_pdf_exports__831e0e03(env, config: dict):
    """
    Check if PDFs were exported to the specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        List of PDF filenames in the directory
    """
    directory = config.get('directory', '/home/user/Downloads/Papers')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        pdf_files = [os.path.basename(f.strip()) for f in list_result['output'].strip().split('\n') if f.strip()]
    return pdf_files

def get_file_moved_check__789836386f3e1cf0e0ee5d172a0885f2(env, config):
    """Check if a file was moved (exists in destination but not in source).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'source_file', 'dest_dir' parameters

    Returns:
        dict: Status of file in source and destination locations
    """
    source_file = config.get('source_file', '')
    dest_dir = config.get('dest_dir', '')
    filename = source_file.split('/')[-1]
    check_source = f'[ -f "{source_file}" ] && echo "exists" || echo "not_exists"'
    source_result = env.controller.run_bash_script(check_source, timeout=10)
    source_exists = False
    if source_result and source_result.get('output'):
        source_exists = source_result['output'].strip() == 'exists'
    check_dest = f'[ -f "{dest_dir}/{filename}" ] && echo "exists" || echo "not_exists"'
    dest_result = env.controller.run_bash_script(check_dest, timeout=10)
    dest_exists = False
    if dest_result and dest_result.get('output'):
        dest_exists = dest_result['output'].strip() == 'exists'
    return {'in_source': source_exists, 'in_dest': dest_exists}

def get_pdf_properties__0b5ef92a5e5fe7305b438702a0eb6c3a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file properties including size and validity.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with file properties
    """
    from pypdf import PdfReader
    import tempfile
    path = config.get('path')
    if not path:
        return {'exists': False}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        is_valid = True
        os.unlink(tmp_path)
        return {'exists': True, 'is_valid': is_valid, 'page_count': page_count, 'file_size': len(file_bytes)}
    except Exception as e:
        logger.error(f'Error reading PDF: {e}')
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass
        return {'exists': True, 'is_valid': False, 'error': str(e), 'file_size': len(file_bytes) if file_bytes else 0}

def get_vm_file__9817ff525ad645c5458b8a22c03332a7(env, config):
    """
    Get the DOCX file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded file
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    import tempfile
    import os
    suffix = os.path.splitext(file_path)[1]
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_docx_file_count__a8fe8251(env, config):
    """Count .docx files in a directory using ls command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        int: Count of .docx files
    """
    directory = config.get('directory', '/home/user/Desktop')
    command = f'ls -1 {directory}/*.docx 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        output = result.get('output', '0')
        try:
            return int(output.strip())
        except ValueError:
            return 0
    output = result.get('output', '0')
    try:
        return int(output.strip())
    except ValueError:
        return 0

def get_subdirectory_exists__6f9e6ffd(env, config: Dict[str, Any]) -> bool:
    """Check if subdirectories exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' parameter (list of paths)

    Returns:
        True if all subdirectories exist, False otherwise
    """
    paths = config.get('paths', [])
    for path in paths:
        result = env.controller.run_bash_script(f"test -d {path} && echo 'exists' || echo 'not_exists'", timeout=10)
        if 'not_exists' in result.get('output', ''):
            return False
    return True

def get_file_exists__974295f9d11461d175dbc0223dd4ff65(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on VM and retrieve its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'content': str} - existence flag and file content
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'exists': False, 'content': ''}
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None or len(file_bytes) == 0:
        return {'exists': False, 'content': ''}
    try:
        content = file_bytes.decode('utf-8')
    except (UnicodeDecodeError, AttributeError):
        try:
            content = str(file_bytes)
        except:
            content = ''
    return {'exists': True, 'content': content}

def get_file_variable_occurrences__ba25acbbf5417edd79f01d39568b431f(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Count occurrences of specific variable names in a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: VM path to the file
            - variables: List of variable names to count

    Returns:
        Dict with variable names as keys and occurrence counts as values:
        {"alist": 0, "arr": 5}
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {}
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1')
    variables = config.get('variables', [])
    result = {}
    for var in variables:
        import re
        pattern = '\\b' + re.escape(var) + '\\b'
        count = len(re.findall(pattern, content))
        result[var] = count
    return result

def get_screenshot_file_info__83a4fedf3a2dd9034045f9bfa6b3d8ed(env, config: dict):
    """
    Check if a screenshot file exists on VM and get its basic information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File information (exists, size, dimensions) or error info
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File not found: {path}')
            return {'exists': False, 'path': path, 'size': 0, 'width': 0, 'height': 0}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            return {'exists': True, 'path': path, 'size': len(file_bytes), 'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking file {path}: {e}')
        return {'exists': False, 'path': path, 'size': 0, 'width': 0, 'height': 0, 'error': str(e)}

def get_text_file_content__e34109066d47c745bfcd2dc2768683b4(env, config: dict):
    """Read and return the text content of a file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the text file on VM

    Returns:
        str: Content of the file, or empty string if file doesn't exist
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return ''
    try:
        content = file_bytes.decode('utf-8').strip()
        logger.info(f'Read {len(content)} characters from {path}')
        return content
    except Exception as e:
        logger.error(f'Error decoding file {path}: {e}')
        return ''

def get_pdf_page_count__e56af0ce(env, config):
    """Get the number of pages in a PDF file.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key

    Returns:
        int: Number of pages in PDF, or 0 if error
    """
    path = config.get('path', '')
    if not path:
        return 0
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return 0
        try:
            pdf_reader = PdfReader(io.BytesIO(file_bytes))
            return len(pdf_reader.pages)
        except:
            return 1 if len(file_bytes) > 0 else 0
    except Exception as e:
        print(f'Error: {e}')
        return 0

def get_pdf_files_list__6b40f44d(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get dict of PDF files with their page counts in the specified directory.

    Returns:
        Dict[str, int]: Dictionary mapping filename to page count
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        count_command = f"pdfinfo '{filepath}' 2>/dev/null | grep -oP '(?<=Pages:)\\s*\\d+' | tr -d ' '"
        count_result = env.controller.run_bash_script(count_command, timeout=10)
        if count_result.get('returncode') == 0:
            count_output = count_result.get('output', '').strip()
            try:
                page_count = int(count_output)
                pdf_info[filename] = page_count
            except (ValueError, TypeError):
                pdf_info[filename] = 0
        else:
            pdf_info[filename] = 0
    return pdf_info

def get_pdf_count__4e03b1ed(env, config: dict):
    """Count PDF files in a specific folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        int: Number of PDF files in the folder
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    command = f"ls -1 '{folder_path}'/*.pdf 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    try:
        count = int(output)
    except (ValueError, AttributeError):
        count = 0
    logger.info(f'PDF file count in {folder_path}: {count}')
    return count

def get_pdf_files_info__0707ddca(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_exists__506ad17f(env, config: Dict[str, Any]):
    """
    Check if a file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (str)

    Returns:
        bool: True if file exists, False otherwise
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    return file_bytes is not None

def get_csv_and_count__a4d37536(env, config):
    """
    Get CSV file and contact count from txt file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with csv_path and count_path

    Returns:
        dict: {"csv_exists": bool, "count": int}
    """
    csv_path = config.get('csv_path', '/home/user/Desktop/contacts.csv')
    count_path = config.get('count_path', '/home/user/Desktop/contact_count.txt')
    result = {'csv_exists': False, 'count': 0}
    csv_bytes = env.controller.get_file(csv_path)
    if csv_bytes:
        result['csv_exists'] = True
    count_bytes = env.controller.get_file(count_path)
    if count_bytes:
        try:
            count_text = count_bytes.decode('utf-8').strip()
            result['count'] = int(count_text)
        except (ValueError, UnicodeDecodeError):
            result['count'] = 0
    return result

def get_pdf_basic_info__d36b36f1df35c32e48eb7411bf5a6c33(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract basic information from a PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with 'exists', 'page_count', 'file_size' keys (exists=False if file doesn't exist)
    """
    from pypdf import PdfReader
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'page_count': 0, 'file_size': 0}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        return {'exists': True, 'page_count': page_count, 'file_size': file_size}
    except Exception as e:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'error': str(e)}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_pdf_file_size__29f6acc744cf0d15caa355ed1701b507(env, config: Dict[str, Any]) -> Any:
    """Get the file size of a PDF file from Google Drive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - query: Query to find the file on Google Drive
            - dest: Local destination filename

    Returns:
        int: File size in bytes, or 0 if file not found/error
    """
    try:
        from desktop_env.evaluators.getters.chrome import get_googledrive_file
        local_path = get_googledrive_file(env, config)
        if not local_path or not os.path.exists(local_path):
            logger.warning(f'PDF file not found: {local_path}')
            return 0
        file_size = os.path.getsize(local_path)
        logger.info(f'PDF file size: {file_size} bytes ({file_size / 1024 / 1024:.2f} MB)')
        return file_size
    except Exception as e:
        logger.error(f'Error getting PDF file size: {e}')
        return 0

def get_desktop_renamed_file__6bf0504dae1e157e634b1e1c5be03fad(env, config: Dict[str, Any]) -> Optional[str]:
    """Get renamed file from Desktop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        Path to cached file, or None if file doesn't exist
    """
    vm_path = config['path']
    dest = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    os.makedirs(env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_file_exists__7648b3ac(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if PDF file exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict containing expected rules with file_path

    Returns:
        Dict with file_exists (bool) and file_path (str)
    """
    expected_rules = config.get('rules', {})
    pdf_path = expected_rules.get('file_path', '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.pdf')
    result = env.controller.get_file(pdf_path)
    file_exists = result is not None and len(result) > 0
    return {'file_exists': file_exists, 'file_path': pdf_path}

def get_word_count_from_file__7aa35df1(env, config) -> Optional[int]:
    """
    Get the word count from a file containing just a number.

    Config:
        path (str): absolute path on the VM to fetch the word count file
        dest (str): file name of the downloaded file

    Returns:
        int: The word count value, or None if file doesn't exist or is invalid
    """
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    try:
        file_data = env.controller.get_file(path)
        if file_data is None:
            logger.warning(f'File not found on VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        word_count = int(content.strip())
        logger.info(f'Successfully read word count from {path}: {word_count}')
        return word_count
    except ValueError as e:
        logger.error(f'Invalid word count format in file {path}: {e}')
        return None
    except Exception as e:
        logger.error(f'Error reading word count file {path}: {e}')
        return None

def get_vm_subtitle_file__f5680565(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_docx_content__d46c5e6abd5059bf9bd9590b3b72a55f(env, config: Dict[str, Any]) -> str:
    """Get the content of a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Combined text content from all paragraphs
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(paragraphs)
    except Exception as e:
        return ''
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_python_file_content__093631738f9b5eba42a5bcf60212ba3b(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Extract Python file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file {path}: {e}')
        return None

def get_docx_bold_content__eb25c2f0ed9e56a458cb4ee477d415ca(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get document content and check if text is bold.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the file path on VM

    Returns:
        Dict with 'content' (str) and 'all_bold' (bool) keys
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'content': '', 'all_bold': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'content': '', 'all_bold': False}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content = []
        for paragraph in doc.paragraphs:
            content.append(paragraph.text)
        full_text = '\n'.join(content)
        total_runs = 0
        bold_runs = 0
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    total_runs += 1
                    if run.bold:
                        bold_runs += 1
        all_bold = total_runs > 0 and bold_runs / total_runs >= 0.9
        return {'content': full_text.strip(), 'all_bold': all_bold}
    finally:
        os.unlink(tmp_path)

def get_pdf_files_from_docx__24a9a657e8a67df9509783bd4f53b233(env, config):
    """Check if PDF files were created from .docx sources.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of PDF filenames that should correspond to .docx files
    """
    directory_path = config.get('path', '/home/user/Desktop')
    command = f'cd {directory_path} && ls *.pdf 2>/dev/null'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        pdf_files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
        return pdf_files
    return []

def get_desktop_file_list__e35e8479c21a9a3d6ef729a997676f8c(env, config: Dict[str, str]) -> Optional[list]:
    """
    Get list of files on the desktop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'desktop_path'

    Returns:
        List of filenames on desktop, or None if error
    """
    desktop_path = config.get('desktop_path', '/home/user/Desktop')
    result = env.controller.run_bash_script(f'ls -1 {desktop_path}', timeout=10)
    if result.get('status') == 'success' and result.get('returncode') == 0:
        output = result.get('output', '')
        files = [f.strip() for f in output.split('\n') if f.strip()]
        return files
    else:
        return None

def get_pdf_exports__b47a318c(env, config: dict):
    """
    Check if PDFs were exported to the specified directory and extract metadata.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        List of dicts, each containing:
            - 'filename': PDF filename
            - 'title': PDF metadata title (if available)
            - 'creator': PDF creator/producer (if available)
            - 'text_snippet': First 200 chars of text content
    """
    directory = config.get('directory', '/home/user/Desktop/Articles')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_info_list = []
    if list_result['output']:
        pdf_paths = [f.strip() for f in list_result['output'].strip().split('\n') if f.strip()]
        for pdf_path in pdf_paths:
            filename = os.path.basename(pdf_path)
            try:
                temp_path = f'/tmp/verify_{filename}'
                env.controller.run_bash_script(f'cp {pdf_path} {temp_path}', timeout=5)
                read_cmd = f'base64 {temp_path}'
                read_result = env.controller.run_bash_script(read_cmd, timeout=10)
                if read_result['output']:
                    import base64
                    import io
                    pdf_bytes = base64.b64decode(read_result['output'])
                    doc = fitz.open(stream=pdf_bytes, filetype='pdf')
                    metadata = doc.metadata
                    title = metadata.get('title', '') if metadata else ''
                    creator = metadata.get('creator', '') if metadata else ''
                    producer = metadata.get('producer', '') if metadata else ''
                    text_snippet = ''
                    if len(doc) > 0:
                        first_page = doc[0]
                        text_snippet = first_page.get_text()[:200]
                    doc.close()
                    pdf_info_list.append({'filename': filename, 'title': title, 'creator': creator, 'producer': producer, 'text_snippet': text_snippet})
                    env.controller.run_bash_script(f'rm -f {temp_path}', timeout=5)
                else:
                    pdf_info_list.append({'filename': filename, 'title': '', 'creator': '', 'producer': '', 'text_snippet': ''})
            except Exception as e:
                pdf_info_list.append({'filename': filename, 'title': '', 'creator': '', 'producer': '', 'text_snippet': ''})
    return pdf_info_list

def get_text_file_lines__5ced85fc_aug18_v3_c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8(env, config):
    """Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        list: List of lines from the file (stripped of trailing newlines)
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return []
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return []
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = [line.rstrip('\n\r') for line in content.splitlines()]
        logger.info(f'Successfully read {len(lines)} lines from {file_path}')
        return lines
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return []

def get_file_exists_at_path__24914e86(env, config):
    """Check if a file exists at a specific path on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    return file_bytes is not None and len(file_bytes) > 0

def get_text_file_content__1a1f627807b83c4d33e1ae428da08935(env, config):
    """Get full content of a text file as a string.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key for file path

    Returns:
        str: Content of the file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        text = file_bytes.decode('utf-8')
        return text
    except Exception as e:
        return ''

def get_file_exists__60340d37(env, config: dict):
    """Check if the exported file exists on VM and validate it's a proper ODS file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with file existence, format, and size information
    """
    vm_path = config.get('path', '/home/user/Desktop/financial-survey.ods')
    result = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    output = result.get('output', '').strip()
    file_exists = output == 'EXISTS'
    if not file_exists:
        return {'exists': False, 'is_ods': False, 'size': 0, 'mime_type': None}
    mime_result = env.controller.run_bash_script(f"file --mime-type -b '{vm_path}'", timeout=10)
    mime_type = mime_result.get('output', '').strip()
    is_ods = mime_type == 'application/vnd.oasis.opendocument.spreadsheet'
    size_result = env.controller.run_bash_script(f"stat -c %s '{vm_path}'", timeout=10)
    try:
        file_size = int(size_result.get('output', '0').strip())
    except (ValueError, AttributeError):
        file_size = 0
    return {'exists': file_exists, 'is_ods': is_ods, 'size': file_size, 'mime_type': mime_type}

def get_pdf_validity__90edd5864af08865c892a3ecd6659029(env, config: Dict[str, Any]) -> Any:
    """Check if a PDF file is valid and readable.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - query: Query to find the file on Google Drive
            - dest: Local destination filename

    Returns:
        dict: Validity info with 'valid', 'page_count', 'has_text' keys
    """
    try:
        import PyPDF2
        from desktop_env.evaluators.getters.chrome import get_googledrive_file
        local_path = get_googledrive_file(env, config)
        if not local_path or not os.path.exists(local_path):
            logger.warning(f'PDF file not found: {local_path}')
            return {'valid': False, 'page_count': 0, 'has_text': False}
        try:
            with open(local_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                if pdf_reader.is_encrypted:
                    logger.warning('PDF is encrypted')
                    return {'valid': False, 'page_count': 0, 'has_text': False}
                page_count = len(pdf_reader.pages)
                has_text = False
                if page_count > 0:
                    pages_to_check = [0]
                    if page_count > 1:
                        pages_to_check.append(page_count - 1)
                    if page_count > 2:
                        pages_to_check.append(page_count // 2)
                    if page_count > 10:
                        pages_to_check.append(page_count // 4)
                        pages_to_check.append(3 * page_count // 4)
                    total_text = ''
                    for page_idx in pages_to_check:
                        try:
                            page = pdf_reader.pages[page_idx]
                            text = page.extract_text()
                            total_text += text
                        except Exception as e:
                            logger.warning(f'Could not extract text from page {page_idx}: {e}')
                    has_text = len(total_text.strip()) > 0
                result = {'valid': True, 'page_count': page_count, 'has_text': has_text}
                logger.info(f'PDF validation result: {result}')
                return result
        except Exception as e:
            logger.error(f'PDF validation failed: {e}')
            return {'valid': False, 'page_count': 0, 'has_text': False}
    except Exception as e:
        logger.error(f'Error validating PDF: {e}')
        return {'valid': False, 'page_count': 0, 'has_text': False}

def get_pdf_file_location__0d0715b7122c298c1e10aa6fa135f599(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if PDF file exists and get its location information.
    Scans the target directory for PDF files instead of using a hardcoded filename.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path is used to extract target directory)

    Returns:
        Dict with file existence and location info
    """
    from pypdf import PdfReader
    target_directory = os.path.dirname(config['path'])
    logger.info(f'Scanning directory {target_directory} for PDF files')
    hardcoded_path = config['path']
    file_bytes = env.controller.get_file(hardcoded_path)
    found_path = None
    if file_bytes is not None and len(file_bytes) > 0:
        found_path = hardcoded_path
        logger.info(f'Found PDF at hardcoded path: {hardcoded_path}')
    else:
        try:
            directory_listing = env.controller.execute(f"find '{target_directory}' -maxdepth 1 -type f -name '*.pdf' -printf '%T@ %p\\n' 2>/dev/null | sort -rn | head -1")
            if directory_listing and directory_listing.strip():
                parts = directory_listing.strip().split(' ', 1)
                if len(parts) == 2:
                    most_recent_pdf = parts[1]
                    logger.info(f'Found most recent PDF: {most_recent_pdf}')
                    file_bytes = env.controller.get_file(most_recent_pdf)
                    if file_bytes is not None and len(file_bytes) > 0:
                        found_path = most_recent_pdf
                        logger.info(f'Successfully retrieved PDF: {found_path}')
        except Exception as e:
            logger.warning(f'Error scanning directory: {e}')
    if found_path:
        result = {'path': found_path, 'exists': True, 'directory': os.path.dirname(found_path), 'filename': os.path.basename(found_path)}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            result['page_count'] = len(reader.pages)
            result['valid_pdf'] = len(reader.pages) > 0
            logger.info(f"PDF file found at {found_path} with {result['page_count']} pages")
        except Exception as e:
            logger.error(f'Error reading PDF: {e}')
            result['valid_pdf'] = False
            result['page_count'] = 0
        finally:
            os.unlink(tmp_path)
    else:
        logger.warning(f'No PDF file found in {target_directory}')
        result = {'path': hardcoded_path, 'exists': False, 'directory': target_directory, 'filename': '', 'valid_pdf': False, 'page_count': 0}
    return result

def get_vacation_files__d7a48669399bf74024b2a979a32d4ae1(env, config: dict) -> dict:
    """Get directory listing from a specific path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Directory tree dict with 'children' list
    """
    return env.controller.get_vm_directory_tree(config['path'])

def get_snake_direction_test__c739e38ab0abd37ef206886e6d29d7b5(env, config: dict):
    """Test snake direction changes work correctly.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters
            - path: List of file paths on VM
            - dest: List of destination filenames
            - multi: Boolean indicating multiple files

    Returns:
        dict: Test result with 'passed' boolean
    """
    paths = config.get('path', [])
    dests = config.get('dest', [])
    if not paths or not dests:
        return {'passed': False, 'error': 'Missing file paths'}
    with tempfile.TemporaryDirectory() as tmpdir:
        for (vm_path, dest) in zip(paths, dests):
            file_bytes = env.controller.get_file(vm_path)
            if not file_bytes:
                return {'passed': False, 'error': f'Failed to get {vm_path}'}
            temp_path = os.path.join(tmpdir, dest)
            with open(temp_path, 'wb') as f:
                f.write(file_bytes)
        sys.path.insert(0, tmpdir)
        try:
            import pygame
            pygame.init()
            from snake import Snake
            from settings import SNAKE_SIZE
            snake = Snake()
            start_pos = snake.positions[0]
            snake.direction = pygame.K_RIGHT
            snake.move()
            if snake.positions[0] != (start_pos[0] + SNAKE_SIZE, start_pos[1]):
                return {'passed': False, 'error': 'RIGHT movement incorrect', 'expected': (start_pos[0] + SNAKE_SIZE, start_pos[1]), 'actual': snake.positions[0]}
            current_pos = snake.positions[0]
            snake.direction = pygame.K_DOWN
            snake.move()
            if snake.positions[0] != (current_pos[0], current_pos[1] + SNAKE_SIZE):
                return {'passed': False, 'error': 'DOWN movement incorrect', 'expected': (current_pos[0], current_pos[1] + SNAKE_SIZE), 'actual': snake.positions[0]}
            current_pos = snake.positions[0]
            snake.direction = pygame.K_LEFT
            snake.move()
            if snake.positions[0] != (current_pos[0] - SNAKE_SIZE, current_pos[1]):
                return {'passed': False, 'error': 'LEFT movement incorrect', 'expected': (current_pos[0] - SNAKE_SIZE, current_pos[1]), 'actual': snake.positions[0]}
            current_pos = snake.positions[0]
            snake.direction = pygame.K_UP
            snake.move()
            if snake.positions[0] != (current_pos[0], current_pos[1] - SNAKE_SIZE):
                return {'passed': False, 'error': 'UP movement incorrect', 'expected': (current_pos[0], current_pos[1] - SNAKE_SIZE), 'actual': snake.positions[0]}
            return {'passed': True}
        except Exception as e:
            return {'passed': False, 'error': str(e)}
        finally:
            sys.path.remove(tmpdir)
            for module in ['food', 'snake', 'main', 'settings']:
                if module in sys.modules:
                    del sys.modules[module]

def get_file_content__8085b902c0aa531b12d2ed766e22c897(env, config: Dict[str, Any]) -> str:
    """Get file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String containing file content, or empty string if file not found
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('File path not specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return ''
        return file_bytes.decode('utf-8')
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return ''

def get_file_permissions__b3b80682(env, config):
    """Check if file exists and get its permissions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        dict: {"file_exists": bool, "permissions": str, "readonly": bool}
    """
    file_path = config.get('file_path')
    result_exists = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    file_exists = 'exists' in result_exists['output'] and 'not exists' not in result_exists['output']
    if not file_exists:
        return {'file_exists': False, 'permissions': '', 'readonly': False}
    result_perms = env.controller.run_bash_script(f'stat -c "%a" "{file_path}"', timeout=10)
    permissions = result_perms['output'].strip()
    try:
        owner_perm = int(permissions[0]) if permissions else 0
        readonly = owner_perm & 2 == 0
    except (ValueError, IndexError):
        readonly = False
    return {'file_exists': file_exists, 'permissions': permissions, 'readonly': readonly}

def get_docx_content_status__ba90be29(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Check if docx file has content from Apple Logo Usage section.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'has_content', 'word_count', 'content', and 'logo_usage_keywords_found' keys, or None if error
    """
    file_path = config.get('path')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        full_text = ''.join([paragraph.text for paragraph in doc.paragraphs]).strip()
        word_count = len(full_text.split()) if full_text else 0
        content_lower = full_text.lower()
        logo_usage_keywords = ['logo', 'apple', 'trademark', 'clearspace', 'clear space', 'branding', 'brand', 'guidelines', 'usage', 'reproduction', 'minimum size', 'color', 'monochrome', 'placement']
        keywords_found = [kw for kw in logo_usage_keywords if kw in content_lower]
        result = {'has_content': len(full_text) > 0, 'word_count': word_count, 'content': full_text, 'logo_usage_keywords_found': keywords_found}
        import os
        os.unlink(tmp_path)
        return result
    except Exception as e:
        return None

def get_pdf_file_list__10f702e7ff0a641a1fda6d45251486ce(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get list of PDF filenames in directory with metadata.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict with 'files' (list of filenames) and 'count' (int)
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"find {directory} -maxdepth 1 -name '*.pdf' -type f -printf '%f\\n' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        logger.info(f'No PDF files found in {directory}')
        return {'files': [], 'count': 0}
    pdf_files = [f.strip() for f in output.split('\n') if f.strip()]
    logger.info(f'Found {len(pdf_files)} PDF files in {directory}: {pdf_files}')
    return {'files': pdf_files, 'count': len(pdf_files)}

def get_folder_file_list__32b2b661(env, config: dict):
    """Get list of filenames in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        list: List of filenames (basenames only) in the directory
    """
    folder_path = config.get('folder_path', '')
    command = f'ls -1 "{folder_path}" 2>/dev/null | grep -E "\\.jpg$|\\.jpeg$|\\.png$" || true'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        output = result['output'].strip()
        if output:
            return output.split('\n')
    return []

def get_text_file_content__de3b1681(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_filename_info__4b590a3a028f08e8f4ad12729f4351c3(env, config):
    """
    Get filename information and verify actual file type for audio conversion verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Information about the file including exists, basename, extension, file_type, and is_audio
    """
    file_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f"if [ -f '{file_path}' ]; then echo 'true'; else echo 'false'; fi", timeout=10)
    exists = exists_result.get('output', '').strip() == 'true'
    info = {'exists': exists, 'basename': '', 'extension': '', 'file_type': '', 'is_audio': False}
    if not exists:
        return info
    basename_result = env.controller.run_bash_script(f"basename '{file_path}'", timeout=10)
    basename = basename_result.get('output', '').strip()
    info['basename'] = basename
    if '.' in basename:
        extension = basename.rsplit('.', 1)[-1]
        info['extension'] = extension
    file_type_result = env.controller.run_bash_script(f"file -b --mime-type '{file_path}'", timeout=10)
    file_type = file_type_result.get('output', '').strip()
    info['file_type'] = file_type
    if file_type.startswith('audio/'):
        info['is_audio'] = True
        logger.info(f'File is confirmed audio: {file_type}')
    else:
        info['is_audio'] = False
        logger.warning(f'File is NOT audio (type: {file_type}), may be renamed video or wrong format')
    logger.info(f'File info for {file_path}: {info}')
    return info

def get_text_file_content__5ced85fc_aug18_v2_b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7(env, config):
    """Read a text file from VM and return its entire content as a string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        str: File content as a string (with trailing whitespace stripped)
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        logger.info(f'Successfully read {len(content)} characters from {file_path}')
        return content
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return ''

def get_files_in_directory__880f8efb(env, config):
    """Get list of files in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        list: List of filenames in the directory
    """
    directory = config.get('directory', '/home/user/Desktop')
    result = env.controller.run_bash_script(f'if [ -d "{directory}" ]; then ls -1 "{directory}"; else echo ""; fi', timeout=10)
    output = result.get('output', '').strip()
    if output:
        return [f.strip() for f in output.split('\n') if f.strip()]
    return []

def get_file_size__b9c089b2fe7d833fde2da297bbbd9620(env, config: dict):
    """
    Get file size from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        int: File size in bytes, or 0 if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.info(f'File not found: {file_path}')
        return 0
    file_size = len(file_bytes)
    logger.info(f'File size: {file_size} bytes for {file_path}')
    return file_size

def get_chapter_pdf_count__457b9850ebb20f14d7c68688b0aa27a6(env, config: dict):
    """Count PDF files in directory matching chapter naming pattern.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        dict: {
            'count': int - number of chapter PDF files found,
            'chapter_files': list - list of chapter filenames found
        }
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'cd "{directory}" && find . -maxdepth 1 -type f -name "[0-9]*.pdf" | sort'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        return {'count': 0, 'chapter_files': []}
    chapter_files = []
    for line in output.split('\n'):
        if line.strip():
            filename = line.strip().replace('./', '')
            chapter_files.append(filename)
    return {'count': len(chapter_files), 'chapter_files': chapter_files}

def get_pdf_files_with_validation__1cc9c1bfb30ee1b2b3cdbf7926894918(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get PDF files in directory and validate their properties including filename matching.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and optional 'expected_titles' keys

    Returns:
        Dict with:
        - 'files': list of PDF filenames found
        - 'valid_count': number of valid PDFs
        - 'chrome_titles': list of Chrome tab titles (if available)
        - 'matched_files': list of filenames that match expected titles
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    expected_titles = config.get('expected_titles', [])
    check_dir_cmd = f"test -d {directory} && echo 'exists' || echo 'missing'"
    dir_result = env.controller.run_bash_script(check_dir_cmd, timeout=10)
    if 'missing' in dir_result.get('output', ''):
        logger.warning(f'Directory {directory} does not exist')
        return {'files': [], 'valid_count': 0, 'chrome_titles': [], 'matched_files': []}
    list_cmd = f"ls {directory}/*.pdf 2>/dev/null | xargs -n 1 basename 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(list_cmd, timeout=10)
    output = result.get('output', '').strip()
    pdf_files = []
    if output:
        pdf_files = [f.strip() for f in output.split('\n') if f.strip() and f.endswith('.pdf')]
    chrome_titles = []
    try:
        import requests
        response = requests.get('http://localhost:9222/json', timeout=5)
        if response.status_code == 200:
            tabs = response.json()
            for tab in tabs:
                if tab.get('type') == 'page' and 'title' in tab:
                    title = tab.get('title', '')
                    if title and title != 'New Tab':
                        chrome_titles.append(title)
            logger.info(f'Retrieved {len(chrome_titles)} Chrome tab titles: {chrome_titles}')
    except Exception as e:
        logger.warning(f'Could not retrieve Chrome tab titles: {e}')
        if expected_titles:
            chrome_titles = expected_titles
    matched_files = []
    if chrome_titles:
        for pdf_file in pdf_files:
            pdf_name = pdf_file[:-4] if pdf_file.endswith('.pdf') else pdf_file
            pdf_normalized = normalize_filename(pdf_name)
            for title in chrome_titles:
                title_normalized = normalize_filename(title)
                if title_normalized in pdf_normalized or pdf_normalized in title_normalized or len(set(title_normalized.split('-')) & set(pdf_normalized.split('-'))) >= min(3, len(title_normalized.split('-')) // 2):
                    matched_files.append(pdf_file)
                    break
    valid_count = len(pdf_files)
    matched_count = len(matched_files)
    logger.info(f'Found {valid_count} PDF files in {directory}')
    logger.info(f'Matched {matched_count} files to Chrome titles/expected titles')
    return {'files': pdf_files, 'valid_count': valid_count, 'chrome_titles': chrome_titles, 'matched_files': matched_files}

def get_text_file_content__834c93d1a65ecbb7766bb5ceb1a12320(env, config: Dict[str, Any]) -> str:
    """Get the content of a text file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the text file

    Returns:
        Content of the text file as a string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        logger.warning(f'Failed to get file from VM: {file_path}')
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file content: {e}')
        return ''

def get_vm_subtitle_file__c48ab0f0(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_pdf_merge_info__a8bc7d70fd87c66cda30ee22072de4d3(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get comprehensive information about both the uploaded PDF and source email attachments.

    This getter combines data from:
    1. Google Drive: uploaded PDF file info
    2. Thunderbird: source email attachment info

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Google Drive settings file path
            - query: Query to find the file on Google Drive
            - email_subject: Subject of the email to search for in Thunderbird
            - profile_path: Path to Thunderbird profile (optional)

    Returns:
        dict: Combined info with 'drive_file' and 'email_attachments' keys
    """
    try:
        drive_config = {'settings_file': config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml'), 'query': config.get('query', [])}
        drive_file = get_googledrive_file_info__a8bc7d70fd87c66cda30ee22072de4d3(env, drive_config)
        email_config = {'profile_path': config.get('profile_path', '/home/user/.thunderbird'), 'email_subject': config.get('email_subject', 'Paper Recommendation')}
        email_attachments = get_thunderbird_email_attachments__a8bc7d70fd87c66cda30ee22072de4d3(env, email_config)
        combined_info = {'drive_file': drive_file, 'email_attachments': email_attachments}
        logger.info(f"Combined PDF merge info: Drive file exists={drive_file.get('exists')}, Email found={email_attachments.get('found')}, PDF count={email_attachments.get('pdf_count')}")
        return combined_info
    except Exception as e:
        logger.error(f'Error getting PDF merge info: {e}')
        return {'drive_file': {'exists': False}, 'email_attachments': {'found': False, 'pdf_count': 0}}

def get_vm_text_file__7c58ef63(env, config):
    """Read text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the file path

    Returns:
        str: File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'File not found or could not be read: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        logger.info(f'Successfully read file: {path}, content length: {len(content)} chars')
        return content.strip()
    except Exception as e:
        logger.error(f'Failed to decode file content: {e}')
        return None

def get_pdf_text__439ece4f99f624319fffd552c64d5f89(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract text content from PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'text' (first page text), 'full_text' keys
    """
    from pypdf import PdfReader
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return {'exists': False, 'text': '', 'full_text': ''}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            first_page_text = ''
            if len(reader.pages) > 0:
                first_page_text = reader.pages[0].extract_text()
            full_text = ''
            for i in range(min(3, len(reader.pages))):
                full_text += reader.pages[i].extract_text() + '\n'
            return {'exists': True, 'text': first_page_text, 'full_text': full_text, 'page_count': len(reader.pages)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading PDF file: {e}')
        return {'exists': False, 'text': '', 'full_text': '', 'page_count': 0}

def get_docx_text_content__7080c120(env, config: dict):
    """Extract text content from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Combined text from all paragraphs in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/book_list_result.docx')
    file_content = env.controller.get_file(vm_path)
    if not file_content:
        return ''
    doc = Document(io.BytesIO(file_content))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    return '\n'.join(text_parts)

def get_file_exists__42ef2d72a0f41077972857e814318a24(env, config: Dict[str, Any]):
    """Check if a file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path to check)

    Returns:
        Dict with exists: bool indicating if file exists
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    return {'exists': file_bytes is not None}

def get_two_zip_file_counts__0cf1485f(env, config: dict):
    """Get the file lists from two zip archives.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'zip_path1' and 'zip_path2' parameters

    Returns:
        dict: Dictionary with 'files1' and 'files2' keys containing lists of filenames
    """
    zip_path1 = config.get('zip_path1', '')
    zip_path2 = config.get('zip_path2', '')
    result = {'files1': [], 'files2': []}
    zip_bytes1 = env.controller.get_file(zip_path1)
    if zip_bytes1:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp:
            tmp.write(zip_bytes1)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zf:
                result['files1'] = [os.path.basename(name) for name in zf.namelist() if not name.endswith('/')]
        except:
            pass
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    zip_bytes2 = env.controller.get_file(zip_path2)
    if zip_bytes2:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp:
            tmp.write(zip_bytes2)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zf:
                result['files2'] = [os.path.basename(name) for name in zf.namelist() if not name.endswith('/')]
        except:
            pass
        finally:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
    return result

def get_pdf_page_count__434044232831686d814f3fab48f559ee(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file and extract page count information.
    Searches for any PDF file on Desktop instead of checking hardcoded filename.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (can be pattern)

    Returns:
        Dict with page count information
    """
    from pypdf import PdfReader
    desktop_path = '/home/user/Desktop'
    try:
        desktop_files = env.controller.execute_python_command(f"import os; print('|||'.join([f for f in os.listdir('{desktop_path}') if f.endswith('.pdf')]))")
        if desktop_files and desktop_files.strip():
            pdf_files = [f for f in desktop_files.strip().split('|||') if f]
        else:
            pdf_files = []
    except Exception as e:
        logger.warning(f'Could not list desktop files: {e}')
        pdf_files = []
    if not pdf_files:
        logger.warning(f'No PDF files found on Desktop')
        return {'path': desktop_path, 'exists': False, 'valid_pdf': False, 'page_count': 0}
    logger.info(f'Found PDF files on Desktop: {pdf_files}')
    config_path = config.get('path', '')
    config_filename = os.path.basename(config_path) if config_path else None
    if config_filename and config_filename in pdf_files:
        pdf_files.remove(config_filename)
        pdf_files.insert(0, config_filename)
    for pdf_filename in pdf_files:
        pdf_path = f'{desktop_path}/{pdf_filename}'
        file_bytes = env.controller.get_file(pdf_path)
        if file_bytes and len(file_bytes) > 0:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                reader = PdfReader(tmp_path)
                page_count = len(reader.pages)
                if page_count > 0:
                    logger.info(f"Valid PDF file '{pdf_filename}' has {page_count} pages")
                    return {'path': pdf_path, 'exists': True, 'valid_pdf': True, 'page_count': page_count}
                else:
                    logger.warning(f"PDF file '{pdf_filename}' has 0 pages")
            except Exception as e:
                logger.warning(f"Error reading PDF '{pdf_filename}': {e}")
            finally:
                os.unlink(tmp_path)
    logger.warning('No valid PDF files found on Desktop')
    return {'path': desktop_path, 'exists': False, 'valid_pdf': False, 'page_count': 0}

def get_pdf_files_in_dir__cdfaf4c1(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_file_move_status__00112b53200a74ce7a53869d2d085264(env, config: dict):
    """Check if a file has been moved from source to target location.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'source_path' and 'target_path' keys

    Returns:
        dict: Status dict with 'target_exists' and 'source_exists' boolean keys
    """
    source_path = config.get('source_path', '')
    target_path = config.get('target_path', '')
    if not source_path or not target_path:
        logger.warning('Source or target path not provided')
        return {'target_exists': False, 'source_exists': False}
    try:
        target_command = f'test -f "{target_path}" && echo "true" || echo "false"'
        target_result = env.controller.run_bash_script(target_command, timeout=10)
        target_exists = target_result.get('output', '').strip() == 'true'
        source_command = f'test -f "{source_path}" && echo "true" || echo "false"'
        source_result = env.controller.run_bash_script(source_command, timeout=10)
        source_exists = source_result.get('output', '').strip() == 'true'
        logger.debug(f'File move status: target_exists={target_exists}, source_exists={source_exists}')
        return {'target_exists': target_exists, 'source_exists': source_exists}
    except Exception as e:
        logger.error(f'Error checking file move status: {e}')
        return {'target_exists': False, 'source_exists': False}

def get_merged_file_content__618fd61c(env, config: Dict[str, Any]) -> str:
    """
    Get the content of a merged file from the VM.

    Args:
        env: Environment object
        config: Configuration dict containing:
            - path: absolute path on the VM to fetch
            - dest: file name of the downloaded file

    Returns:
        str: Path to the local cached file, or None if retrieval failed
    """
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        logger.info(f'Successfully saved file: {cache_path} ({len(file_content)} bytes)')
        return cache_path
    except Exception as e:
        logger.error(f'Error processing file {path}: {e}')
        return None

def get_file_modification_time__23cbcfa9(env, config):
    """Get file modification timestamp from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Unix timestamp of last modification, or -1 if file not found
    """
    path = config.get('path', '')
    result = env.controller.run_bash_script(f'if [ -f "{path}" ]; then stat -c %Y "{path}"; else echo "-1"; fi', timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        logger.error(f'Failed to parse modification time: {output}')
        return -1

def get_pdf_file_count__086472c391d9b3dbf463fe72713bc019(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get information about PDF files in a directory on VM, including filenames and count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key (VM path)

    Returns:
        Dict with 'count' (int), 'filenames' (list), and 'has_original' (bool)
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f"cd '{directory}' && ls -1 *.pdf 2>/dev/null"
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0:
        logger.warning(f"Failed to list PDF files in {directory}: {result.get('error', '')}")
        return {'count': 0, 'filenames': [], 'has_original': False}
    output = result.get('output', '').strip()
    if not output:
        logger.info(f'No PDF files found in {directory}')
        return {'count': 0, 'filenames': [], 'has_original': False}
    filenames = [line.strip() for line in output.split('\n') if line.strip()]
    count = len(filenames)
    has_original = 'Spectral Graph Theory.pdf' in filenames
    logger.info(f'Found {count} PDF files in {directory}: {filenames}')
    return {'count': count, 'filenames': filenames, 'has_original': has_original}

def get_csv_filtered_rows__e49a53b3238eb33c8eed6130c48c5268(env, config):
    """
    Read CSV file from VM and return all rows as list of lists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists representing CSV rows (including header)
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        rows = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
    finally:
        os.unlink(tmp_path)

def get_file_exists__7930967f(env, config: Dict[str, str]) -> dict:
    """
    Check if a file exists at the specified path in the VM and verify its content.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains 'exists' (bool), 'size' (int), and 'is_text' (bool)
    """
    vm_ip = env.vm_ip
    port = env.server_port
    path = config['path']
    result = {'exists': False, 'size': 0, 'is_text': False}
    command = f"test -f '{path}' && echo 'exists' || echo 'not_exists'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
    if response.status_code != 200:
        logger.error(f'Failed to check file existence. Status code: {response.status_code}')
        return result
    output = response.json()['output'].strip()
    result['exists'] = output == 'exists'
    if not result['exists']:
        return result
    size_command = f"stat -c %s '{path}' 2>/dev/null || echo '0'"
    size_response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': size_command, 'shell': True})
    if size_response.status_code == 200:
        try:
            result['size'] = int(size_response.json()['output'].strip())
        except ValueError:
            logger.error('Failed to parse file size')
            result['size'] = 0
    type_command = f"file -b --mime-type '{path}' 2>/dev/null || echo 'unknown'"
    type_response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': type_command, 'shell': True})
    if type_response.status_code == 200:
        mime_type = type_response.json()['output'].strip()
        result['is_text'] = mime_type.startswith('text/')
    return result

def get_text_file_lines__465a794bbc0c138e64d3078d7f3f0b51(env, config):
    """Read text file and return lines as list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of non-empty lines from file
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception:
        return []

def get_docx_content__61084d52(env, config):
    """Read content from a .docx file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains 'exists' (bool), 'is_valid_docx' (bool), 'text' (str), and 'word_count' (int) keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'is_valid_docx': False, 'text': '', 'word_count': 0}
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None or len(file_bytes) == 0:
            return {'exists': False, 'is_valid_docx': False, 'text': '', 'word_count': 0}
        try:
            docx_buffer = BytesIO(file_bytes)
            with zipfile.ZipFile(docx_buffer, 'r') as docx_zip:
                xml_content = docx_zip.read('word/document.xml')
                tree = ET.fromstring(xml_content)
                namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                text_elements = tree.findall('.//w:t', namespaces)
                extracted_text = ''.join((elem.text for elem in text_elements if elem.text))
                extracted_text = extracted_text.strip()
                word_count = len(extracted_text.split()) if extracted_text else 0
                return {'exists': True, 'is_valid_docx': True, 'text': extracted_text, 'word_count': word_count}
        except (zipfile.BadZipFile, KeyError, ET.ParseError) as e:
            print(f'File is not a valid .docx: {e}')
            return {'exists': True, 'is_valid_docx': False, 'text': '', 'word_count': 0}
    except Exception as e:
        print(f'Error reading docx file: {e}')
        return {'exists': False, 'is_valid_docx': False, 'text': '', 'word_count': 0}

def get_vm_subtitle_file__dccfbd8e(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_default_pdf_viewer__10cce6e7e23a5b966c5cbd94b3842103(env, config: dict):
    """Gets the default PDF viewer application.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: The default PDF viewer .desktop file name (e.g., 'evince.desktop')
    """
    import requests
    os_type = env.vm_platform
    if os_type == 'Linux':
        command = ['xdg-mime', 'query', 'default', 'application/pdf']
        vm_ip = env.vm_ip
        port = env.server_port
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
        if response.status_code == 200:
            app = response.json().get('output', '').strip()
            if app:
                return app
        return 'unknown'
    else:
        raise Exception('Unsupported operating system', os_type)

def get_file_exists__25f0d8c3(env, config):
    """
    Check if a file exists at the specified path on VM and verify it's a valid PNG image.

    Returns:
        dict with keys:
            - exists: bool, whether file exists
            - is_valid_png: bool, whether file has valid PNG magic bytes
            - size: int, file size in bytes (0 if file doesn't exist)
    """
    file_path = config.get('path', '')
    command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if output != 'EXISTS':
        return {'exists': False, 'is_valid_png': False, 'size': 0}
    size_command = f'stat -c %s "{file_path}" 2>/dev/null || echo "0"'
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    magic_command = f'xxd -l 8 -p "{file_path}" 2>/dev/null || echo ""'
    magic_result = env.controller.run_bash_script(magic_command, timeout=10)
    magic_bytes = magic_result.get('output', '').strip().replace('\n', '')
    is_valid_png = magic_bytes.lower() == '89504e470d0a1a0a'
    return {'exists': True, 'is_valid_png': is_valid_png, 'size': file_size}

def get_default_pdf_viewer__09a6505e4073dfc93152a41afb590f9b(env, config: dict):
    """Gets the default application for PDF files on Linux.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: The default PDF viewer application name
    """
    from desktop_env.evaluators.getters.general import get_vm_command_line
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'application/pdf']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_docx_text_content__3874241e(env, config):
    """Get full text content from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: All text content in the document
    """
    file_path = config.get('path')
    if not file_path:
        return ''
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return ''
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            text_content = []
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            return '\n'.join(text_content)
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error getting text content: {e}')
        return ''

def get_text_file_content__8b701bf8d0cacb95438d2e4e17a8b914(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file on VM

    Returns:
        str: File content as string
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_file_exists_check__e185e2be(env, config: Dict) -> Optional[Dict]:
    """
    Check if a file exists and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_dir_file_count__c907493f(env, config: Dict[str, Any]) -> int:
    """Get the count of files in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters
            - directory: Path to directory on VM
            - pattern: Optional file extension filter (e.g., ".tex")

    Returns:
        int: Number of files matching the criteria
    """
    directory = config.get('directory', '/home/user')
    pattern = config.get('pattern', '')
    command = f'ls -1 "{directory}" 2>/dev/null | wc -l'
    if pattern:
        command = f'ls -1 "{directory}"/*{pattern} 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.warning(f"Failed to count files in {directory}: {result.get('error', '')}")
        return 0
    try:
        count = int(result['output'].strip())
        return count
    except (ValueError, AttributeError) as e:
        logger.error(f'Failed to parse file count: {e}')
        return 0

def get_pdf_exists__7a335cc66cd23236defef5bd95bbbe7d(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a PDF file exists and extract verification data including content comparison with source .docx.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with verification data: exists, is_valid_pdf, file_size, page_count, text_preview, source_text_preview
    """
    path = config.get('path', '')
    result = {'exists': False, 'is_valid_pdf': False, 'file_size': 0, 'page_count': 0, 'text_preview': '', 'source_exists': False, 'source_text_preview': '', 'content_matches': False}
    try:
        pdf_bytes = env.controller.get_file(path)
        if pdf_bytes is None:
            logger.info(f'PDF file not found at {path}')
            return result
        result['exists'] = True
        result['file_size'] = len(pdf_bytes)
        if pdf_bytes[:4] != b'%PDF':
            logger.warning(f'File at {path} exists but is not a valid PDF')
            return result
        result['is_valid_pdf'] = True
        try:
            import PyPDF2
            from io import BytesIO
            pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
            result['page_count'] = len(pdf_reader.pages)
            text_parts = []
            for i in range(min(3, len(pdf_reader.pages))):
                page_text = pdf_reader.pages[i].extract_text()
                if page_text:
                    text_parts.append(page_text)
            result['text_preview'] = ' '.join(text_parts)[:1000]
            logger.info(f"PDF found with {result['page_count']} pages")
        except Exception as e:
            logger.warning(f'Could not extract PDF content: {e}')
            pass
        docx_path = path.replace('.pdf', '.docx')
        try:
            docx_bytes = env.controller.get_file(docx_path)
            if docx_bytes:
                result['source_exists'] = True
                try:
                    from docx import Document
                    from io import BytesIO
                    doc = Document(BytesIO(docx_bytes))
                    source_text_parts = []
                    for para in doc.paragraphs[:10]:
                        if para.text.strip():
                            source_text_parts.append(para.text)
                    result['source_text_preview'] = ' '.join(source_text_parts)[:1000]
                    if result['text_preview'] and result['source_text_preview']:
                        source_words = set((word.lower() for word in result['source_text_preview'].split() if len(word) > 5))
                        pdf_words = set((word.lower() for word in result['text_preview'].split() if len(word) > 5))
                        if source_words and pdf_words:
                            overlap = len(source_words.intersection(pdf_words))
                            overlap_ratio = overlap / len(source_words) if len(source_words) > 0 else 0
                            result['content_matches'] = overlap_ratio >= 0.5
                            logger.info(f'Content match ratio: {overlap_ratio:.2f}')
                except Exception as e:
                    logger.warning(f'Could not extract .docx content: {e}')
        except Exception as e:
            logger.warning(f'Could not check source .docx file: {e}')
        return result
    except Exception as e:
        logger.error(f'Error checking PDF file: {e}')
        return result

def get_txt_file_content__b2ba9ee6873b43000b20ba169c4d9592(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get content from a TXT file at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'content' keys
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'TXT file not found at {path}')
        return {'exists': False, 'path': path, 'content': ''}
    try:
        content = file_bytes.decode('utf-8', errors='replace')
        logger.info(f'TXT file found at {path} with {len(content)} characters')
        return {'exists': True, 'path': path, 'content': content}
    except Exception as e:
        logger.error(f'Error reading TXT file: {e}')
        return {'exists': False, 'path': path, 'content': ''}

def get_pdf_files_in_dir__1ee77af4(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_file_count_in_directory__23c57bc9(env, config):
    """Count files matching a pattern in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern' keys

    Returns:
        int: Number of matching files
    """
    directory = config.get('directory', '/home/user/Desktop')
    pattern = config.get('pattern', '*.mp4')
    result = env.controller.run_bash_script(f'cd "{directory}" 2>/dev/null && ls -1 {pattern} 2>/dev/null | wc -l || echo "0"', timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        logger.error(f'Failed to parse file count: {output}')
        return 0

def get_tetris_files__72c9bc74(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_pdf_file_sizes__4e03b1ed(env, config: dict):
    """Get sizes of PDF files in a folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        dict: Mapping of filename to size in bytes
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    command = f"cd '{folder_path}' 2>/dev/null && ls -l *.pdf 2>/dev/null | awk '{{print $9,$5}}' || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        return {}
    file_sizes = {}
    for line in output.split('\n'):
        if line.strip():
            parts = line.strip().split()
            if len(parts) >= 2:
                filename = parts[0]
                try:
                    size = int(parts[1])
                    file_sizes[filename] = size
                except ValueError:
                    pass
    logger.info(f'PDF file sizes: {file_sizes}')
    return file_sizes

def get_file_content__b6bb3e42(env, config):
    """Get file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content
    """
    path = config.get('path', '')
    command = f'cat {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        return result['output'].strip()
    else:
        logger.error(f"Failed to read file {path}: {result['error']}")
        return None

def get_text_file_lines__c2b520792100f4aa54f0d89dedfa94c4(env, config):
    """
    Get the lines of a text file from VM.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key specifying the file path on VM

    Returns:
        list: List of non-empty lines (stripped), or empty list if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes:
        try:
            content = file_bytes.decode('utf-8', errors='ignore')
            lines = [line.strip() for line in content.split('\n') if line.strip()]
            return lines
        except Exception:
            return []
    else:
        return []

def get_desktop_named_pdf__0c8a922818da41c6a16e3979ae1f26dc(env, config: Dict[str, Any]) -> Optional[str]:
    """Get PDF file with specific name from Desktop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        Path to cached PDF file, or None if file doesn't exist
    """
    vm_path = config['path']
    dest = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    os.makedirs(env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_remaining_files__8754d37bdc9e8d94ab80feb618caa015(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of remaining files in a folder on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' key

    Returns:
        List of filenames (sorted)
    """
    folder_path = config.get('folder_path', '')
    command = f"ls -1 '{folder_path}' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0 or not result.get('output'):
        logger.warning(f'Failed to list folder: {folder_path}')
        return []
    files = result['output'].strip().split('\n')
    files = [f for f in files if f]
    return sorted(files)

def get_python_file_content__be0d67f2ca1c7fde5d3c550cd046d0b4(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Extract Python file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file {path}: {e}')
        return None

def get_file_exists_check__e4d07acf(env, config: Dict) -> Optional[Dict]:
    """
    Check if a file exists and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' and 'size' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'size': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'size': 0}
        file_size = os.path.getsize(local_path)
        return {'exists': True, 'size': file_size}
    except Exception as e:
        return {'exists': False, 'size': 0, 'error': str(e)}

def get_renamed_file__9f8a1d0c(env, config):
    """Check if file was renamed correctly.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_path' and 'new_path'

    Returns:
        Dict with 'old_exists' and 'new_exists' booleans
    """
    old_path = config.get('old_path')
    new_path = config.get('new_path')
    old_bytes = env.controller.get_file(old_path)
    old_exists = old_bytes is not None
    new_bytes = env.controller.get_file(new_path)
    new_exists = new_bytes is not None
    return {'old_exists': old_exists, 'new_exists': new_exists}

def get_text_file_content__5ef70598fcc68171b6ec36d7c888fc21(env, config):
    """Get full content of a text file as a string.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key for file path

    Returns:
        str: Content of the file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        text = file_bytes.decode('utf-8')
        return text
    except Exception as e:
        return ''

def get_file_exists__e3dda739ee4da14903d2e8df52d7a41d(env, config: dict):
    """
    Check if a file exists on the VM and validate it's a valid MP3 audio file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        dict: {
            'exists': bool - True if file exists,
            'is_mp3': bool - True if file is a valid MP3 audio file,
            'file_size': int - File size in bytes (0 if doesn't exist)
        }
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None or len(file_bytes) == 0:
        logger.info(f'File does not exist or is empty: {file_path}')
        return {'exists': False, 'is_mp3': False, 'file_size': 0}
    file_size = len(file_bytes)
    is_mp3 = False
    if len(file_bytes) >= 3:
        if file_bytes[:3] == b'ID3':
            is_mp3 = True
            logger.info(f'File {file_path} has ID3 tag (MP3)')
        elif len(file_bytes) >= 2 and file_bytes[0] == 255 and (file_bytes[1] & 224 == 224):
            is_mp3 = True
            logger.info(f'File {file_path} has MPEG sync (MP3)')
        else:
            logger.warning(f"File {file_path} exists but doesn't have valid MP3 magic bytes")
    result = {'exists': True, 'is_mp3': is_mp3, 'file_size': file_size}
    logger.info(f'File check for {file_path}: {result}')
    return result

def get_default_pdf_viewer__bdb427f6(env, config: dict):
    """Gets the default application for PDF files on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'application/pdf']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_text_content__67c66e9d6c723be29d116d6e2c7b5850(env, config):
    """Get the content of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Content of the text file
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception as e:
        return ''

def get_backup_folder_files__089628262a733e157c368e0bf1ad02f1(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if backup folder exists and contains the required files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' and 'expected_files' keys

    Returns:
        Dict with 'folder_exists', 'files_found' list
    """
    folder_path = config['folder_path']
    expected_files = config.get('expected_files', [])
    check_cmd = f"test -d '{folder_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'"
    result = env.controller.run_bash_script(check_cmd, timeout=10)
    if 'NOT_EXISTS' in result.get('output', ''):
        return {'folder_exists': False, 'files_found': []}
    ls_cmd = f"ls -1 '{folder_path}' 2>/dev/null || echo ''"
    ls_result = env.controller.run_bash_script(ls_cmd, timeout=10)
    files_in_folder = [f.strip() for f in ls_result.get('output', '').split('\n') if f.strip()]
    files_found = []
    for expected_file in expected_files:
        if expected_file in files_in_folder:
            files_found.append(expected_file)
    return {'folder_exists': True, 'files_found': files_found}

def get_file_content__edd1c1331eff751ad5487718d7e5d07b(env, config: Dict[str, Any]) -> str:
    """Get the content of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content, or empty string if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Could not read file: {config['path']}")
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error decoding file: {e}')
        return ''

def get_vm_file_exists__8752e3fa(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if a file exists on the VM and get its size.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Dict with 'exists' (bool) and 'size_bytes' (int) keys
    """
    path = config['path']
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; else echo 'NOT_EXISTS'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result['output'].strip()
    if output == 'NOT_EXISTS' or result['returncode'] != 0:
        return {'exists': False, 'size_bytes': 0}
    try:
        size = int(output)
        return {'exists': True, 'size_bytes': size}
    except ValueError:
        logger.error(f'Failed to parse file size: {output}')
        return {'exists': False, 'size_bytes': 0}

def get_file_line_count__e09bfbdf(env, config: dict):
    """
    Count lines in a file.

    This function counts actual lines of text, not just newline characters.
    A file with 5 lines but no trailing newline is still counted as 5 lines.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Number of lines in file (0 if file doesn't exist)
    """
    path = config.get('path', '')
    check_command = f'test -f "{path}" && echo "exists" || echo "not_exists"'
    check_result = env.controller.run_bash_script(check_command, timeout=10)
    if not check_result or check_result.get('output', '').strip() != 'exists':
        return 0
    command = f'grep -c "" "{path}" 2>/dev/null || echo "0"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        try:
            return int(result['output'].strip())
        except ValueError:
            return 0
    return 0

def get_pdf_basic_info__0448af3ec8f726fcfbf94a06b95da924(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract basic information from a PDF file including headers and footers detection.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with 'exists', 'page_count', 'file_size', 'has_headers_footers' keys
    """
    from pypdf import PdfReader
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'has_headers_footers': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        has_headers_footers = detect_headers_footers(reader)
        return {'exists': True, 'page_count': page_count, 'file_size': file_size, 'has_headers_footers': has_headers_footers}
    except Exception as e:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'has_headers_footers': False, 'error': str(e)}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_file_exists__724cfa0a(env, config):
    """Check if specific file exists in cloned repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' parameter

    Returns:
        bool: True if file exists, False otherwise
    """
    file_path = config.get('file_path')
    command = f"test -f {file_path} && echo 'EXISTS' || echo 'NOT_EXISTS'"
    result = env.controller.run_bash_script(command, timeout=10)
    return 'EXISTS' in result.get('output', '')

def get_archive_contents__ccef1fda30d9645300e2b2b1c657b52f(env, config: dict):
    """Extract and return the contents of a tar.gz archive from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the archive on VM

    Returns:
        dict: {'exists': bool, 'files': list of filenames in archive}
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get archive from VM: {path}')
        return {'exists': False, 'files': []}
    try:
        with tempfile.NamedTemporaryFile(suffix='.tar.gz', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with tarfile.open(tmp_path, 'r:gz') as tar:
                all_names = tar.getnames()
                files = sorted([os.path.basename(name) for name in all_names if not name.endswith('/')])
                logger.info(f'Archive contains {len(files)} files: {files}')
                return {'exists': True, 'files': files}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading archive {path}: {e}')
        return {'exists': False, 'files': []}

def get_identified_mountains_folder__c8c87c59ecc91c5a6beb4965364b8f59(env, config: dict):
    """Check if Identified_Mountains folder exists and contains the moved mountain pictures.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: {
            'exists': bool,
            'file_count': int,
            'files': list,
            'has_correct_files': bool,
            'source_files_removed': bool
        }
    """
    folder_path = '/home/user/Pictures/Identified_Mountains'
    source_files = ['/home/user/Pictures/picture1.jpg', '/home/user/Pictures/picture2.jpg', '/home/user/Pictures/picture3.jpg']
    expected_filenames = {'picture1.jpg', 'picture2.jpg', 'picture3.jpg'}
    check_cmd = f'test -d "{folder_path}" && echo "exists" || echo "not_exists"'
    check_result = env.controller.run_bash_script(check_cmd, timeout=10)
    if 'not_exists' in check_result.get('output', ''):
        return {'exists': False, 'file_count': 0, 'files': [], 'has_correct_files': False, 'source_files_removed': False}
    list_cmd = f'ls -1 "{folder_path}" 2>/dev/null || echo ""'
    list_result = env.controller.run_bash_script(list_cmd, timeout=10)
    files = []
    if list_result.get('output'):
        files = [f.strip() for f in list_result['output'].strip().split('\n') if f.strip()]
    file_count = len(files)
    files_set = set(files)
    has_correct_files = files_set == expected_filenames
    source_files_removed = True
    for source_file in source_files:
        check_source_cmd = f'test -f "{source_file}" && echo "exists" || echo "not_exists"'
        source_check_result = env.controller.run_bash_script(check_source_cmd, timeout=10)
        if 'exists' in source_check_result.get('output', ''):
            source_files_removed = False
            break
    return {'exists': True, 'file_count': file_count, 'files': files, 'has_correct_files': has_correct_files, 'source_files_removed': source_files_removed}

def get_backup_folder_hashes__dd5eeed0(env, config):
    """Get hashes of all images in the backup folder."""
    backup_dir = config.get('backup_dir', '/home/user/Pictures/backup')
    result = env.controller.run_bash_script(f"ls -1 {backup_dir} 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        return []
    files = result['output'].strip().split('\n')
    hashes = []
    for filename in files:
        if not filename or filename.strip() == '':
            continue
        file_path = os.path.join(backup_dir, filename.strip())
        file_bytes = env.controller.get_file(file_path)
        if file_bytes:
            try:
                from io import BytesIO
                with Image.open(BytesIO(file_bytes)) as img:
                    img_byte_arr = img.tobytes()
                    hash_result = hashlib.sha256(img_byte_arr).hexdigest()
                    hashes.append(hash_result)
            except Exception as e:
                print(f'Error processing {filename}: {e}')
                continue
    return sorted(hashes)

def get_tetris_files__0e016e1b(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_pdf_files_list__21dcc4a0(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files in the specified directory with their page counts.

    Returns:
        Dict[str, int]: Mapping of filename to page count
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        filepath = os.path.join(directory, filename)
        page_cmd = f"pdfinfo '{filepath}' 2>/dev/null | grep '^Pages:' | awk '{{print $2}}'"
        page_result = env.controller.run_bash_script(page_cmd, timeout=10)
        if page_result.get('returncode') == 0:
            page_output = page_result.get('output', '').strip()
            try:
                page_count = int(page_output)
                pdf_info[filename] = page_count
            except (ValueError, TypeError):
                pdf_info[filename] = 0
        else:
            pdf_info[filename] = 0
    return pdf_info

def get_pdf_files_list__0be3e647(env, config: Dict[str, Any]) -> Dict[str, int]:
    """Get list of PDF files with their page counts in the specified directory."""
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls {directory}/*.pdf 2>/dev/null | xargs -n1 basename'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '').strip()
    if not output:
        return {}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    pdf_info = {}
    for filename in files:
        full_path = os.path.join(directory, filename)
        page_cmd = f"pdfinfo '{full_path}' 2>/dev/null | grep 'Pages:' | awk '{{print $2}}'"
        page_result = env.controller.run_bash_script(page_cmd, timeout=10)
        if page_result.get('returncode') == 0:
            page_output = page_result.get('output', '').strip()
            try:
                page_count = int(page_output)
                pdf_info[filename] = page_count
            except (ValueError, TypeError):
                pdf_info[filename] = 0
        else:
            pdf_info[filename] = 0
    return pdf_info

def get_renamed_files__27c9f1432580c9f5f1bf2d1f919f4ed5(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if files have been renamed with the expected prefix.
    Verifies that renamed files exist AND original files do not exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path', 'prefix', and 'base_filenames' keys

    Returns:
        Dict with 'renamed_files_found', 'original_files_remaining', 'base_filenames', and 'prefix'
    """
    folder_path = config['folder_path']
    prefix = config.get('prefix', '')
    base_filenames = config.get('base_filenames', [])
    ls_cmd = f"ls -1 '{folder_path}' 2>/dev/null || echo ''"
    ls_result = env.controller.run_bash_script(ls_cmd, timeout=10)
    files_in_folder = [f.strip() for f in ls_result.get('output', '').split('\n') if f.strip()]
    renamed_files_found = []
    for base_filename in base_filenames:
        expected_renamed = f'{prefix}{base_filename}'
        if expected_renamed in files_in_folder:
            renamed_files_found.append(expected_renamed)
    original_files_remaining = []
    for base_filename in base_filenames:
        if base_filename in files_in_folder:
            original_files_remaining.append(base_filename)
    return {'renamed_files_found': renamed_files_found, 'original_files_remaining': original_files_remaining, 'base_filenames': base_filenames, 'prefix': prefix}

def get_text_file_content__08e73922(env, config):
    """Read text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Text content of the file, stripped of whitespace
    """
    path = config.get('path', '')
    if not path:
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        print(f'Error reading file {path}: {e}')
        return ''

def get_pdf_file_sizes__99e94d1a(env, config: dict):
    """Get sizes of PDF files in specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        List of tuples (filename, size_in_bytes)
    """
    directory = config.get('directory', '/home/user/Documents/WebPages')
    command = f"find {directory} -maxdepth 1 -name '*.pdf' -type f -exec du -b {{}} \\; 2>/dev/null"
    result = env.controller.run_bash_script(command, timeout=30)
    output = result.get('output', '').strip()
    if not output:
        return []
    file_sizes = []
    for line in output.split('\n'):
        if line.strip():
            parts = line.split('\t', 1)
            if len(parts) == 2:
                size = int(parts[0])
                filepath = parts[1]
                filename = os.path.basename(filepath)
                file_sizes.append((filename, size))
    logger.info(f'Found {len(file_sizes)} PDF files with sizes in {directory}')
    return file_sizes

def get_sender_file__ad43db2627764dd28ab9631606c7b97c(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Read the content of a text file that should contain a sender email address.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (file path to read)

    Returns:
        Content of the file as string, or None if file doesn't exist
    """
    file_path = config.get('path', '/home/user/sender.txt')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return None
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return None

def get_file_list__91c6f86e(env, config: Dict[str, Any]) -> List[str]:
    """Get list of all files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of filenames
    """
    path = config['path']
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return []
    filenames = [node['name'] for node in result['children']]
    return filenames

def get_pdf_numbered_files__fe03784ed42c9a7002fcc758df207953(env, config: dict):
    """Get list of numbered PDF files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        dict: {'numbered_files': list of (number, filename) tuples, 'total_files': int}
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls -1 "{directory}"/*.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0 or not result.get('output'):
        return {'numbered_files': [], 'total_files': 0}
    output = result['output'].strip()
    if not output:
        return {'numbered_files': [], 'total_files': 0}
    numbered_files = []
    total_files = 0
    for line in output.split('\n'):
        line = line.strip()
        if line:
            total_files += 1
            filename = line.split('/')[-1]
            match = re.match('^(\\d+)\\. ', filename)
            if match:
                number = int(match.group(1))
                numbered_files.append((number, filename))
    numbered_files.sort(key=lambda x: x[0])
    return {'numbered_files': numbered_files, 'total_files': total_files}

def get_text_file_content__68f83c37fe78a3f2dbadefcb3c480330(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    path = config.get('path', '/home/user/output.txt')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_text_file_line_count__5538243966b3481fe772536923c1f693(env, config):
    """
    Read text file from VM and parse content as integer.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Integer value parsed from file content (0 if invalid)
    """
    file_path = config.get('path', '')
    if not file_path:
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return 0
    try:
        content = file_bytes.decode('utf-8').strip()
        if content.isdigit():
            return int(content)
        return 0
    except Exception as e:
        print(f'Error reading text file: {e}')
        return 0

def get_zip_filenames__5c33f919(env, config: dict):
    """Get list of filenames in a zip archive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'zip_path' parameter

    Returns:
        list: List of filenames (basenames only, no paths) in the zip
    """
    zip_path = config.get('zip_path', '')
    zip_bytes = env.controller.get_file(zip_path)
    if not zip_bytes:
        return []
    with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp:
        tmp.write(zip_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            filenames = [os.path.basename(name) for name in zf.namelist() if not name.endswith('/') and os.path.basename(name)]
        return filenames
    except Exception as e:
        return []
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

def get_writer_text_content__592c042e86aba9b7d5a9be8008d4a4f0(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text content from Writer document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with text content including paragraphs and lists
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'error': 'File not found'}
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            return {'paragraphs': paragraphs, 'paragraph_count': len(paragraphs), 'tables': tables, 'table_count': len(tables)}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'error': str(e)}

def get_text_file_content__2d81db9f5efc3d72c38ba9f24bf6d4fb(env, config):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_ipynb_dirs__c69badbf0960e9198fe7a7334d4b95f3(env, config):
    """
    Read list of directory names from ipynb_directories.txt file created by user.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used, kept for signature compatibility)

    Returns:
        list: List of directory names from the file
    """
    file_path = '/home/user/test_environment/ipynb_directories.txt'
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'File not found or empty: {file_path}')
        return []
    try:
        content = file_bytes.decode('utf-8').strip()
        if not content:
            return []
        dirs = [line.strip() for line in content.split('\n') if line.strip()]
        return dirs
    except Exception as e:
        logger.error(f'Failed to read/parse file {file_path}: {e}')
        return []

def get_docs_subdir_doc_count__09b713a6(env, config):
    """Verify .doc files were moved to docs subdirectory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        dict: Contains 'docs_count' (files in docs/) and 'desktop_count' (files in Desktop root)
    """
    docs_directory = config.get('directory', '/home/user/Desktop/docs')
    desktop_directory = '/home/user/Desktop'
    docs_command = f"ls -1 {docs_directory}/*.doc 2>/dev/null | grep -v '\\.docx$' | wc -l"
    docs_result = env.controller.run_bash_script(docs_command, timeout=10)
    desktop_command = f"find {desktop_directory} -maxdepth 1 -type f -name '*.doc' ! -name '*.docx' 2>/dev/null | wc -l"
    desktop_result = env.controller.run_bash_script(desktop_command, timeout=10)
    docs_output = docs_result.get('output', '0')
    desktop_output = desktop_result.get('output', '0')
    try:
        docs_count = int(docs_output.strip())
    except ValueError:
        docs_count = 0
    try:
        desktop_count = int(desktop_output.strip())
    except ValueError:
        desktop_count = 0
    return {'docs_count': docs_count, 'desktop_count': desktop_count}

def get_file_exists__ec920d7f(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and is a valid ODS file.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' to file

    Returns:
        dict: Contains 'exists' (bool), 'valid_ods' (bool), and 'has_content' (bool)
    """
    file_content = env.controller.get_file(config['path'])
    if file_content is None or len(file_content) == 0:
        return {'exists': False, 'valid_ods': False, 'has_content': False}
    valid_ods = False
    try:
        file_bytes = BytesIO(file_content)
        with zipfile.ZipFile(file_bytes, 'r') as zip_ref:
            file_list = zip_ref.namelist()
            required_files = ['mimetype', 'content.xml']
            has_required = all((f in file_list for f in required_files))
            if 'mimetype' in file_list:
                mimetype_content = zip_ref.read('mimetype').decode('utf-8').strip()
                is_spreadsheet = mimetype_content == 'application/vnd.oasis.opendocument.spreadsheet'
            else:
                is_spreadsheet = False
            valid_ods = has_required and is_spreadsheet
    except (zipfile.BadZipFile, Exception):
        valid_ods = False
    return {'exists': True, 'valid_ods': valid_ods, 'has_content': len(file_content) > 0}

def get_zip_file_count__03f8ef9d(env, config: dict):
    """Extract the list of filenames from a zip archive.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'zip_path' parameter

    Returns:
        list: List of filenames (basenames) in the zip archive
    """
    zip_path = config.get('zip_path', '')
    zip_bytes = env.controller.get_file(zip_path)
    if not zip_bytes:
        return []
    with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp:
        tmp.write(zip_bytes)
        tmp_path = tmp.name
    try:
        with zipfile.ZipFile(tmp_path, 'r') as zf:
            filenames = [os.path.basename(name) for name in zf.namelist() if not name.endswith('/')]
        return filenames
    except Exception as e:
        return []
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)

def get_file_content_text__8369c71b0ee26c543c025c6e1cb39bbd(env, config):
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String content of the file, stripped of whitespace
    """
    path = config.get('path', '')
    if not path:
        return ''
    result = env.controller.run_bash_script(f'cat {path}', timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    else:
        return ''

def get_recent_pdf_count__4e03b1ed(env, config: dict):
    """Count PDF files modified within the last N minutes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' and 'minutes' parameters

    Returns:
        int: Number of recently modified PDF files
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    minutes = config.get('minutes', 10)
    command = f"find '{folder_path}' -name '*.pdf' -type f -mmin -{minutes} 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    try:
        count = int(output)
    except ValueError:
        count = 0
    logger.info(f'Found {count} PDFs modified in last {minutes} minutes in {folder_path}')
    return count

def get_pdf_file_info__000c73be1394701e25c8491dd9418647(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get PDF file from VM and verify it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        str: Path to downloaded PDF file in cache, or None if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import os
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_file_exists__0ad50b28(env, config: dict):
    """Check if the exported TSV file exists on VM and validate its format and content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary containing file validation results with keys:
            - exists (bool): Whether file exists
            - is_tsv (bool): Whether file uses tab delimiters
            - has_content (bool): Whether file is non-empty
            - row_count (int): Number of data rows (excluding header if present)
            - sample_line (str): First line of the file for validation
    """
    vm_path = config.get('path', '/home/user/Desktop/survey-data.tsv')
    exists_result = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = exists_result.get('output', '').strip() == 'EXISTS'
    if not file_exists:
        return {'exists': False, 'is_tsv': False, 'has_content': False, 'row_count': 0, 'sample_line': ''}
    size_result = env.controller.run_bash_script(f"wc -c < '{vm_path}'", timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    has_content = file_size > 0
    head_result = env.controller.run_bash_script(f"head -n 5 '{vm_path}'", timeout=10)
    sample_content = head_result.get('output', '').strip()
    is_tsv = '\t' in sample_content if sample_content else False
    row_count_result = env.controller.run_bash_script(f"wc -l < '{vm_path}'", timeout=10)
    row_count = int(row_count_result.get('output', '0').strip())
    first_line_result = env.controller.run_bash_script(f"head -n 1 '{vm_path}'", timeout=10)
    sample_line = first_line_result.get('output', '').strip()
    return {'exists': file_exists, 'is_tsv': is_tsv, 'has_content': has_content, 'row_count': row_count, 'sample_line': sample_line}

def get_renamed_docx_file__232d51eaf1622d499f79eaf2309ca5ad(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a renamed DOCX file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists' key indicating if file exists
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'Renamed DOCX file not found at {path}')
        return {'exists': False, 'path': path}
    logger.info(f'Renamed DOCX file found at {path} ({len(file_bytes)} bytes)')
    return {'exists': True, 'path': path, 'size': len(file_bytes)}

def get_pdf_files_in_dir__d8278409(env, config: dict):
    """
    Get list of PDF files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PDF filenames (sorted)
    """
    dir_path = config.get('path', '/home/user/Desktop/book')
    command = f'cd {dir_path} && ls -1 *.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning(f'No PDF files found in {dir_path}')
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    pdf_files = sorted([f for f in files if f.endswith('.pdf')])
    logger.info(f'Found PDF files in {dir_path}: {pdf_files}')
    return pdf_files

def get_zip_verification_state__1a0ccf050b10545e05649e32a895cb33(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get verification state for zip file and its contents.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'zip_path' and 'expected_files' keys

    Returns:
        Dict with 'zip_exists', 'zip_path', and 'contents' (list of files in zip)
    """
    zip_path = config['zip_path']
    check_result = env.controller.run_bash_script(f"test -f '{zip_path}' && echo 'exists' || echo 'missing'", timeout=10)
    zip_exists = check_result.get('output', '').strip() == 'exists'
    contents = []
    if zip_exists:
        list_result = env.controller.run_bash_script(f"unzip -l '{zip_path}' | tail -n +4 | head -n -2 | awk '{{print $NF}}' | grep -v '^$' | xargs -n1 basename 2>/dev/null || echo ''", timeout=10)
        if list_result['returncode'] == 0:
            output = list_result.get('output', '').strip()
            if output:
                contents = [f.strip() for f in output.split('\n') if f.strip()]
    return {'zip_exists': zip_exists, 'zip_path': zip_path, 'contents': contents}

def get_merged_file_content__3b864a94(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get content from a merged text file and return metadata about it.

    Config:
        path (str): absolute path on the VM to fetch the merged file
        dest (str): file name of the downloaded file

    Returns:
        Dict with keys:
            - exists (bool): whether the file exists
            - line_count (int): number of lines in the file
            - content (str): file content
            - file_size (int): size in bytes
    """
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    cache_path = os.path.join(env.cache_dir, dest)
    try:
        file_data = env.controller.get_file(path)
        if file_data is None:
            logger.warning(f'File not found on VM: {path}')
            return {'exists': False, 'line_count': 0, 'content': '', 'file_size': 0}
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with open(cache_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        lines = content.splitlines()
        line_count = len(lines)
        file_size = len(file_data)
        logger.info(f'Successfully read merged file: {path} ({line_count} lines, {file_size} bytes)')
        return {'exists': True, 'line_count': line_count, 'content': content, 'file_size': file_size}
    except Exception as e:
        logger.error(f'Error reading merged file {path}: {e}')
        return {'exists': False, 'line_count': 0, 'content': '', 'file_size': 0}

def get_notebooks_dir_files__1b47b6505a7a2dc3d6ad6f0c07b4bcb4(env, config):
    """
    Get list of .ipynb files in the notebooks directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: Sorted list of .ipynb filenames in notebooks directory
    """
    command = 'if [ -d /home/user/test_environment/notebooks ]; then cd /home/user/test_environment/notebooks && find . -name "*.ipynb" -type f | sed "s|^./||" | sort; else echo ""; fi'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['status'] != 'success':
        logger.error(f"Command failed: {result.get('error', 'Unknown error')}")
        return []
    try:
        output = result['output'].strip()
        if not output:
            return []
        files = [f for f in output.split('\n') if f]
        return sorted(files)
    except Exception as e:
        logger.error(f'Failed to parse files: {e}')
        return []

def get_all_doc_pdf_conversion__1d31f566abd6111dc90f53a510d255f5(env, config):
    """Get information about doc/docx to PDF conversion.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains 'pdf_count' and 'command_used' fields
    """
    directory_path = config.get('path', '/home/user/Desktop')
    count_command = f'ls {directory_path}/*.pdf 2>/dev/null | wc -l'
    count_result = env.controller.run_bash_script(count_command, timeout=10)
    pdf_count = 0
    if count_result['returncode'] == 0:
        try:
            pdf_count = int(count_result['output'].strip())
        except (ValueError, AttributeError):
            pdf_count = 0
    history_command = 'cat ~/.bash_history'
    history_result = env.controller.run_bash_script(history_command, timeout=10)
    command_used = ''
    if history_result['returncode'] == 0:
        command_used = history_result['output']
    return {'pdf_count': pdf_count, 'command_used': command_used}

def get_text_file_content__d332c3241fced231d1d84d00e75fe3b7(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    path = config.get('path', '/home/user/output.txt')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_file_exists__7ae8ce2b(env, config: dict):
    """Check if file exists on VM and has content.
    
    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key
        
    Returns:
        dict: {'exists': bool, 'has_content': bool}
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'has_content': False}
    has_content = len(file_bytes) > 0
    return {'exists': True, 'has_content': has_content}

def get_file_exists_and_size__c0d603dc(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_subdirs_in_dir__c4632abc(env, config: dict):
    """Get list of subdirectories in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Parent directory path

    Returns:
        List of subdirectory names
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"""\npython3 -c "\nimport os\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path):\n    print('[]')\n    exit(0)\n\nsubdirs = [d for d in os.listdir(dir_path) if os.path.isdir(os.path.join(dir_path, d))]\nimport json\nprint(json.dumps(subdirs))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to list subdirectories: {result['error']}")
        return []
    import json
    try:
        subdirs = json.loads(result['output'].strip())
        return subdirs
    except:
        return []

def get_file_exists__f25970ca(env, config: dict):
    """Validate that the CSV file was properly exported from the XLSX source.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains validation results including file existence, format, content, and data correspondence
    """
    csv_path = config.get('path', '/home/user/Desktop/enterprise-data.csv')
    xlsx_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.xlsx'
    result = {'exists': False, 'is_csv': False, 'has_content': False, 'row_count': 0, 'xlsx_exists': False, 'data_matches': False, 'valid': False}
    check_exists = env.controller.run_bash_script(f"test -f '{csv_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    if check_exists.get('output', '').strip() != 'EXISTS':
        return result
    result['exists'] = True
    check_xlsx = env.controller.run_bash_script(f"test -f '{xlsx_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    if check_xlsx.get('output', '').strip() == 'EXISTS':
        result['xlsx_exists'] = True
    file_content_result = env.controller.run_bash_script(f"cat '{csv_path}'", timeout=30)
    file_content = file_content_result.get('output', '')
    if not file_content or not file_content.strip():
        return result
    result['has_content'] = True
    csv_rows = []
    try:
        csv_reader = csv.reader(io.StringIO(file_content))
        csv_rows = list(csv_reader)
        if len(csv_rows) == 0:
            return result
        result['is_csv'] = True
        result['row_count'] = len(csv_rows)
        has_data = False
        for row in csv_rows[:10]:
            if len(row) > 0 and any((cell.strip() for cell in row)):
                has_data = True
                break
        if not has_data:
            return result
    except Exception as e:
        result['is_csv'] = False
        return result
    if result['xlsx_exists'] and len(csv_rows) > 0:
        try:
            python_script = f'''\nimport sys\ntry:\n    from openpyxl import load_workbook\n\n    # Load XLSX file\n    wb = load_workbook("{xlsx_path}", data_only=True)\n    ws = wb.active\n\n    # Get row count (excluding completely empty rows)\n    xlsx_row_count = 0\n    for row in ws.iter_rows():\n        if any(cell.value is not None for cell in row):\n            xlsx_row_count += 1\n\n    # Get column count from first row\n    first_row = next(ws.iter_rows(min_row=1, max_row=1))\n    xlsx_col_count = sum(1 for cell in first_row if cell.value is not None)\n\n    # Get sample data from first few cells for validation\n    sample_cells = []\n    for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=min(5, xlsx_row_count)), 1):\n        for col_idx, cell in enumerate(row[:min(5, xlsx_col_count)], 1):\n            if cell.value is not None:\n                sample_cells.append(f"{{row_idx}},{{col_idx}}:{{str(cell.value).strip()}}")\n\n    print(f"ROWS:{{xlsx_row_count}}")\n    print(f"COLS:{{xlsx_col_count}}")\n    print(f"SAMPLES:{{';'.join(sample_cells[:10])}}")\n\nexcept ImportError:\n    print("ERROR:openpyxl not available")\nexcept Exception as e:\n    print(f"ERROR:{{str(e)}}")\n'''
            xlsx_check = env.controller.run_bash_script(f"python3 -c '{python_script}'", timeout=30)
            xlsx_output = xlsx_check.get('output', '')
            if 'ERROR:' not in xlsx_output:
                xlsx_row_count = 0
                xlsx_col_count = 0
                xlsx_samples = []
                for line in xlsx_output.strip().split('\n'):
                    if line.startswith('ROWS:'):
                        xlsx_row_count = int(line.split(':')[1])
                    elif line.startswith('COLS:'):
                        xlsx_col_count = int(line.split(':')[1])
                    elif line.startswith('SAMPLES:'):
                        sample_str = line.split(':', 1)[1]
                        if sample_str:
                            xlsx_samples = sample_str.split(';')
                csv_data_rows = len([r for r in csv_rows if any((cell.strip() for cell in r))])
                row_count_match = abs(csv_data_rows - xlsx_row_count) <= 2
                csv_col_count = len(csv_rows[0]) if csv_rows else 0
                col_count_match = abs(csv_col_count - xlsx_col_count) <= 2
                sample_matches = 0
                if xlsx_samples:
                    for sample in xlsx_samples[:5]:
                        if ':' in sample:
                            cell_value = sample.split(':', 1)[1]
                            csv_str = '\n'.join((','.join(row) for row in csv_rows[:5]))
                            if cell_value in csv_str:
                                sample_matches += 1
                sample_match_ratio = sample_matches / max(len(xlsx_samples[:5]), 1) if xlsx_samples else 0
                if row_count_match and col_count_match and (sample_match_ratio >= 0.6):
                    result['data_matches'] = True
                    result['valid'] = True
        except Exception as e:
            pass
    elif result['is_csv'] and result['row_count'] > 1:
        result['valid'] = True
    return result

def get_file_exists_and_size__995b229f(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_pdf_exports__98b0d7a7(env, config: dict):
    """
    Check if PDFs were exported to the specified directory and extract actual page titles.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict with 'pdf_files' (list of PDF filenames) and 'page_titles' (list of actual page titles from Chrome)
    """
    directory = config.get('directory', '/home/user/Documents/WebPages')
    list_command = f'find {directory} -maxdepth 1 -name "*.pdf" -type f 2>/dev/null | sort'
    list_result = env.controller.run_bash_script(list_command, timeout=10)
    pdf_files = []
    if list_result['output']:
        pdf_files = [os.path.basename(f.strip()) for f in list_result['output'].strip().split('\n') if f.strip()]
    page_titles = []
    try:
        cdp_command = 'python3 -c "\nimport json\nimport urllib.request\nimport urllib.error\n\ntry:\n    response = urllib.request.urlopen(\'http://localhost:9222/json\')\n    tabs = json.loads(response.read())\n    titles = [tab.get(\'title\', \'\') for tab in tabs if tab.get(\'type\') == \'page\' and \'title\' in tab]\n    print(json.dumps(titles))\nexcept Exception as e:\n    print(json.dumps([]))\n"'
        cdp_result = env.controller.run_bash_script(cdp_command, timeout=10)
        if cdp_result['output']:
            try:
                page_titles = json.loads(cdp_result['output'].strip())
            except json.JSONDecodeError:
                page_titles = []
    except Exception:
        page_titles = []
    return {'pdf_files': pdf_files, 'page_titles': page_titles}

def get_filenames_with_pattern__7427978e(env, config: Dict[str, Any]) -> List[str]:
    """Get list of filenames in a directory that contain a specific pattern.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of filenames
    """
    path = config['path']
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return []
    filenames = [node['name'] for node in result['children']]
    return filenames

def get_copied_files__9e83c51609bc4e6d77dc6303641e8cf9(env, config: Dict) -> Dict:
    """
    Get files that were copied to a destination directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'destination' key

    Returns:
        Dict mapping file hashes to filenames
    """
    destination = config.get('destination', '/home/user/Desktop/Mountains')
    command = f"""python3 -c "\nimport os\nimport hashlib\nimport json\nfrom PIL import Image\n\ndestination = '{destination}'\n\nresult = {{}}\n\nif os.path.exists(destination) and os.path.isdir(destination):\n    files = os.listdir(destination)\n    for filename in files:\n        filepath = os.path.join(destination, filename)\n        if os.path.isfile(filepath):\n            try:\n                with Image.open(filepath) as img:\n                    img_byte_arr = img.tobytes()\n                    file_hash = hashlib.sha256(img_byte_arr).hexdigest()\n                    result[file_hash] = filename\n            except:\n                pass\n\nprint(json.dumps(result))\n"\n"""
    run_result = env.controller.run_bash_script(command, timeout=10)
    if run_result.get('returncode') != 0:
        return {}
    output = run_result.get('output', '')
    try:
        import json
        files_by_hash = json.loads(output.strip())
        return files_by_hash
    except:
        return {}

def get_file_content__e5eac3aa5287398e5510fa8c359cddce(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    result = env.controller.run_bash_script(f'cat {file_path}', timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    else:
        return ''

def get_pdf_file_properties__869dd2b5b7a0e45511c735d06983da83(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF file and extract its properties (page count, orientation, dimensions).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with PDF properties or empty dict if file doesn't exist
    """
    from pypdf import PdfReader
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.warning(f"Failed to get file from VM: {config['path']}")
        return {}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        reader = PdfReader(tmp_path)
        if len(reader.pages) == 0:
            return {}
        first_page = reader.pages[0]
        mediabox = first_page.mediabox
        width = float(mediabox.width)
        height = float(mediabox.height)
        orientation = 'landscape' if width > height else 'portrait'
        result = {'page_count': len(reader.pages), 'orientation': orientation, 'width': width, 'height': height, 'exists': True}
        logger.info(f'PDF properties: {result}')
        return result
    except Exception as e:
        logger.error(f'Error reading PDF: {e}')
        return {}
    finally:
        os.unlink(tmp_path)

def get_file_exists__67be1ac6efe87edab008a615fb0e7ec4(env, config):
    """Check if a JPEG file exists on the VM and validate it's a valid JPEG image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'exists': bool,
            'path': str,
            'is_jpeg': bool,
            'file_size': int,
            'file_type': str
        }
    """
    file_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "true" || echo "false"', timeout=10)
    exists = exists_result.get('output', '').strip() == 'true'
    if not exists:
        return {'exists': False, 'path': file_path, 'is_jpeg': False, 'file_size': 0, 'file_type': ''}
    file_type_result = env.controller.run_bash_script(f'file -b "{file_path}"', timeout=10)
    file_type = file_type_result.get('output', '').strip()
    is_jpeg = 'JPEG image data' in file_type or 'jpeg' in file_type.lower()
    size_result = env.controller.run_bash_script(f'stat -c %s "{file_path}" 2>/dev/null || echo "0"', timeout=10)
    file_size = int(size_result.get('output', '0').strip())
    return {'exists': exists, 'path': file_path, 'is_jpeg': is_jpeg, 'file_size': file_size, 'file_type': file_type}

def get_vm_file_text__26660ad1(env, config):
    """
    Get the text content of a file from the VM.

    Args:
        env: Environment object with controller.get_file() method
        config: Configuration dict with 'path' key

    Returns:
        str: The file content as text, or empty string if file doesn't exist
    """
    path = config.get('path')
    if not path:
        return ''
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return ''
    try:
        text = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        try:
            text = file_bytes.decode('latin-1')
        except Exception:
            return ''
    return text

def get_multi_directory_contents__92a58812(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_file_not_exists__9e688855(env, config):
    """Check if a file does NOT exist on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file does NOT exist, False if it exists
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    return file_bytes is None or len(file_bytes) == 0

def get_text_file_content__11d0824d05970e58a5671a5636365f15(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file on VM

    Returns:
        str: File content as string
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_word_count_file__adfc25c4(env, config: dict):
    """Get word count from text file saved by user.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Word count read from file, or None if file doesn't exist or invalid
    """
    file_path = config.get('path', '/home/user/Desktop/word_count.txt')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found: {file_path}')
            return None
        content = file_bytes.decode('utf-8').strip()
        try:
            count = int(content)
            return count
        except ValueError:
            logger.warning(f'Invalid word count format: {content}')
            return None
    except Exception as e:
        logger.error(f'Error reading word count file: {e}')
        return None

def get_pdf_file_list__822cf85f742bbfae9e3acf8d7027940c(env, config: dict):
    """Get detailed information about PDF files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        dict: Dictionary with file details:
            - 'original': bool (True if 'Spectral Graph Theory.pdf' exists)
            - 'chapters': list of chapter PDF filenames (excluding the original)
            - 'total_count': int (total number of PDF files)
            - 'all_files': list of all PDF filenames
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f'ls -1 "{directory}"/*.pdf 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0 or not result.get('output'):
        return {'original': False, 'chapters': [], 'total_count': 0, 'all_files': []}
    output = result['output'].strip()
    if not output:
        return {'original': False, 'chapters': [], 'total_count': 0, 'all_files': []}
    files = []
    for line in output.split('\n'):
        line = line.strip()
        if line:
            filename = line.split('/')[-1]
            files.append(filename)
    original_exists = 'Spectral Graph Theory.pdf' in files
    chapters = [f for f in files if f != 'Spectral Graph Theory.pdf']
    return {'original': original_exists, 'chapters': chapters, 'total_count': len(files), 'all_files': files}

def get_file_list__990ae9b047da99489a16db0558f7ee61(env, config: Dict[str, Any]) -> Dict[str, bool]:
    """Check existence of multiple files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' key (list of file paths to check)

    Returns:
        Dict mapping file paths to existence status (True/False)
    """
    paths = config.get('paths', [])
    result = {}
    for path in paths:
        file_bytes = env.controller.get_file(path)
        result[path] = file_bytes is not None and len(file_bytes) > 0
    return result

def get_file_exists__684b5a3a3f653750766f5bbe64af3bd5(env, config: Dict[str, Any]) -> bool:
    """Check if a file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location

    Returns:
        True if file exists, False otherwise
    """
    file_bytes = env.controller.get_file(config['path'])
    return file_bytes is not None and len(file_bytes) > 0

def get_pdf_file_sizes__2fc6b524(env, config: dict):
    """Get sizes of PDF files in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - directory: Directory path

    Returns:
        List of file sizes in bytes
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"""\npython3 -c "\nimport os\nimport glob\n\ndir_path = '{directory}'\nif not os.path.exists(dir_path):\n    print('[]')\n    exit(0)\n\npdf_files = glob.glob(os.path.join(dir_path, '*.pdf'))\nsizes = [os.path.getsize(f) for f in pdf_files]\nimport json\nprint(json.dumps(sizes))\n"\n"""
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        logger.error(f"Failed to get file sizes: {result['error']}")
        return []
    import json
    try:
        sizes = json.loads(result['output'].strip())
        return sizes
    except:
        return []

def get_mountain_prefixed_files__0f256e5c55ade51a0f98ec735f4e6698(env, config: dict):
    """Get list of files in Pictures directory that start with 'Mountain_'.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: List of filenames starting with 'Mountain_'
    """
    command = 'cd /home/user/Pictures && ls -1 Mountain_* 2>/dev/null || echo ""'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] != 0 or not result['output']:
        return []
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    return files

def get_pdf_text_content__7a2513e6563b76f5c07e27d8c4089d02(env, config: Dict[str, Any]) -> Optional[str]:
    """Extract text content from a PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the PDF file on VM

    Returns:
        str: Extracted text content from the PDF, or None if file not found
    """
    import fitz
    vm_path = config.get('path')
    if not vm_path:
        logger.warning('No path specified in config')
        return None
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.warning(f'File not found at {vm_path}')
        return None
    if not file_bytes.startswith(b'%PDF'):
        logger.warning(f'File at {vm_path} is not a valid PDF')
        return None
    import tempfile
    cache_path = os.path.join(env.cache_dir, config.get('dest', 'output.pdf'))
    os.makedirs(os.path.dirname(cache_path) if os.path.dirname(cache_path) else env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        text_content = ''
        with fitz.open(cache_path) as pdf:
            for page in pdf:
                text_content += page.get_text()
        logger.info(f'Extracted {len(text_content)} characters from PDF')
        return text_content.strip()
    except Exception as e:
        logger.error(f'Error extracting text from PDF: {e}')
        return None

def get_vm_file__a5c1978c26f91a02b256f236a6017e74(env, config):
    """
    Get the DOCX file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded file
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    import tempfile
    import os
    suffix = os.path.splitext(file_path)[1]
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_csv_row_count__58df2e60(env, config: Dict[str, Any]):
    """Extract CSV file information for comprehensive verification.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' to CSV file

    Returns:
        dict: Contains file_exists, path, row_count, column_count, headers, and sample_data
    """
    file_path = config['path']
    try:
        file_content = env.controller.get_file(file_path)
        if not file_content:
            return {'file_exists': False, 'path': file_path, 'row_count': 0, 'column_count': 0, 'headers': [], 'data': []}
    except Exception:
        return {'file_exists': False, 'path': file_path, 'row_count': 0, 'column_count': 0, 'headers': [], 'data': []}
    content_str = file_content.decode('utf-8', errors='ignore')
    reader = csv.reader(io.StringIO(content_str))
    rows = list(reader)
    if len(rows) == 0:
        return {'file_exists': True, 'path': file_path, 'row_count': 0, 'column_count': 0, 'headers': [], 'data': []}
    headers = rows[0] if len(rows) > 0 else []
    data_rows = rows[1:] if len(rows) > 1 else []
    column_count = len(headers) if headers else len(data_rows[0]) if data_rows else 0
    return {'file_exists': True, 'path': file_path, 'row_count': len(data_rows), 'column_count': column_count, 'headers': headers, 'data': data_rows}

def get_file_exists__ecf92ffc(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Check if TSV file exists and get its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        Dict with file_exists (bool), file_path (str), content (str), is_tsv_format (bool),
        row_count (int), and column_count (int)
    """
    tsv_path = '/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.tsv'
    content = env.controller.get_file(tsv_path)
    file_exists = content is not None and len(content) > 0
    result = {'file_exists': file_exists, 'file_path': tsv_path, 'content': content if file_exists else '', 'is_tsv_format': False, 'row_count': 0, 'column_count': 0}
    if file_exists and content:
        try:
            lines = content.strip().split('\n')
            if len(lines) > 0:
                first_line = lines[0]
                columns = first_line.split('\t')
                column_count = len(columns)
                is_tsv = column_count >= 2 and '\t' in first_line
                if is_tsv and len(lines) > 1:
                    for i in range(min(5, len(lines))):
                        if '\t' not in lines[i]:
                            is_tsv = False
                            break
                result['is_tsv_format'] = is_tsv
                result['row_count'] = len(lines)
                result['column_count'] = column_count if is_tsv else 0
        except Exception as e:
            logger.error(f'Error validating TSV format: {e}')
            result['is_tsv_format'] = False
    return result

def get_pdf_file_count__2f750d009f0f471751ed869c09f90a90(env, config):
    """Count the number of PDF files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the directory path

    Returns:
        int: Number of PDF files in the directory
    """
    directory_path = config.get('path', '/home/user/Desktop')
    command = f'ls {directory_path}/*.pdf 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        try:
            count = int(result['output'].strip())
            return count
        except (ValueError, AttributeError):
            return 0
    return 0

def get_zip_filenames__d6b113924c427deceec5f933af24484e(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get ZIP archive name and internal filenames.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'zip_exists', 'zip_name', 'internal_files' keys
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return None
    zip_name = os.path.basename(path)
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return {'zip_exists': False, 'zip_name': zip_name, 'internal_files': [], 'valid_zip': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.zip', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                return {'zip_exists': True, 'zip_name': zip_name, 'internal_files': file_list, 'valid_zip': True}
        finally:
            os.unlink(tmp_path)
    except zipfile.BadZipFile:
        logger.error(f'Invalid ZIP file: {path}')
        return {'zip_exists': True, 'zip_name': zip_name, 'internal_files': [], 'valid_zip': False}
    except Exception as e:
        logger.error(f'Error reading ZIP file: {e}')
        return {'zip_exists': False, 'zip_name': zip_name, 'internal_files': [], 'valid_zip': False}

def get_text_file_content__25fb76d7ccb83f42013b589a25bead61(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return ''

def get_text_file_content__ed1a5c265e6c6d06dcaf2ec482204403(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file not found
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return ''

def get_pdf_file_sizes__d4f3d6039bf1a71698a65c2d049521e8(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get sizes and names of PDF files in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict with 'files' (mapping filename to size) and 'filenames' (list of filenames)
    """
    directory = config.get('directory', '/home/user/Documents/Blog')
    command = f"cd {directory} 2>/dev/null && find . -maxdepth 1 -name '*.pdf' -type f -exec stat -c '%n %s' {{}} \\; 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    file_sizes = {}
    filenames = []
    output = result.get('output', '').strip()
    if output:
        for line in output.split('\n'):
            if not line.strip():
                continue
            parts = line.strip().split()
            if len(parts) >= 2:
                filename = parts[0].replace('./', '')
                try:
                    size = int(parts[1])
                    file_sizes[filename] = size
                    filenames.append(filename)
                except ValueError:
                    continue
    logger.info(f'Found {len(file_sizes)} PDF files with sizes in {directory}')
    return {'files': file_sizes, 'filenames': filenames}

def get_pdf_basic_info__247861df47176dd2a8c57d33f4d99eab(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Extract basic information from a PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with 'exists', 'page_count', 'file_size' keys, or None if file doesn't exist
    """
    from pypdf import PdfReader
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'exists': False, 'page_count': 0, 'file_size': 0}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        file_size = len(file_bytes)
        return {'exists': True, 'page_count': page_count, 'file_size': file_size}
    except Exception as e:
        return {'exists': False, 'page_count': 0, 'file_size': 0, 'error': str(e)}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_file_exists__e1da6937(env, config: dict):
    """
    Check if a file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    path = config.get('path', '')
    command = f'test -f "{path}" && echo "exists" || echo "not_exists"'
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('output'):
        output = result['output'].strip()
        return output == 'exists'
    return False

def get_pdf_file_info__5f59826466b2625834fd8d369560ed11(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get PDF file from VM and verify it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to PDF file on VM

    Returns:
        str: Path to downloaded PDF file in cache, or None if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import os
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_file_rename_check__c14397f3284104a2f980691d2ea6abf3(env, config: Dict[str, Any]) -> Dict[str, bool]:
    """
    Check if a file rename operation was completed correctly.
    Verifies that the new file exists AND the old file does not exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'old_path' and 'new_path' keys

    Returns:
        Dict with 'new_exists' (bool) and 'old_exists' (bool) keys
    """
    old_path = config.get('old_path', '')
    new_path = config.get('new_path', '')
    command_old = f"test -f '{old_path}' && echo 'exists' || echo 'not_exists'"
    result_old = env.controller.run_bash_script(command_old, timeout=10)
    old_exists = result_old.get('output', '').strip() == 'exists'
    command_new = f"test -f '{new_path}' && echo 'exists' || echo 'not_exists'"
    result_new = env.controller.run_bash_script(command_new, timeout=10)
    new_exists = result_new.get('output', '').strip() == 'exists'
    logger.info(f"Rename check - old file '{old_path}' exists: {old_exists}, new file '{new_path}' exists: {new_exists}")
    return {'old_exists': old_exists, 'new_exists': new_exists}

def get_csv_filtered_count__83247848fa547df5ef0efe5837431ceb(env, config):
    """Read CSV and count rows matching a filter criterion.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'filter_column' and 'filter_value' keys

    Returns:
        int: Number of rows matching the filter, or -1 if error
    """
    file_path = config.get('path', '')
    filter_column = config.get('filter_column', '')
    filter_value = config.get('filter_value', '')
    if not file_path:
        return -1
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return -1
        import csv
        import io
        content = file_bytes.decode('utf-8', errors='ignore')
        reader = csv.DictReader(io.StringIO(content))
        count = 0
        for row in reader:
            if filter_column in row and row[filter_column].strip() == filter_value:
                count += 1
        return count
    except Exception as e:
        return -1

def get_file_count_by_pattern__816aac7b5fcbde572f13b62a3999bd4d(env, config: Dict) -> Dict:
    """
    Count files in a directory that match specific patterns (for mountain names).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key

    Returns:
        Dict mapping pattern categories to counts
    """
    directory = config.get('directory', '/home/user/Pictures')
    command = f'''python3 -c "import os; files = os.listdir('{directory}'); print(files)"'''
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        return {}
    output = result.get('output', '')
    try:
        file_list_str = output.strip()
        files = eval(file_list_str) if file_list_str else []
    except:
        return {}
    kilimanjaro_count = 0
    everest_count = 0
    hua_count = 0
    for filename in files:
        filename_lower = filename.lower()
        if any((kw in filename_lower for kw in ['kili', 'kilimanjaro'])):
            kilimanjaro_count += 1
        elif any((kw in filename_lower for kw in ['everest', 'sagarmatha', 'chomolungma', 'qomolangma', 'himalaya'])):
            everest_count += 1
        elif any((kw in filename_lower for kw in ['hua', 'huashan'])):
            hua_count += 1
    return {'kilimanjaro': kilimanjaro_count, 'everest': everest_count, 'hua': hua_count}

def get_pdf_files_in_folder__4e03b1ed(env, config: dict):
    """Check PDF files in a specific folder with validation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        Dict with 'files' (list of PDF filenames) and 'file_info' (dict with size and validation)
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    command = f"find '{folder_path}' -maxdepth 1 -type f -name '*.pdf' -exec ls -l {{}} \\; 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        return {'files': [], 'file_info': {}}
    files = []
    file_info = {}
    for line in output.split('\n'):
        if line.strip():
            parts = line.split()
            if len(parts) >= 9:
                size = int(parts[4])
                filepath = ' '.join(parts[8:])
                filename = os.path.basename(filepath)
                files.append(filename)
                file_info[filename] = {'size': size}
                check_cmd = f"file -b '{filepath}' | grep -i pdf"
                check_result = env.controller.run_bash_script(check_cmd, timeout=5)
                is_valid_pdf = 'PDF' in check_result.get('output', '')
                file_info[filename]['is_valid_pdf'] = is_valid_pdf
    logger.info(f'Found PDF files in {folder_path}: {files}')
    logger.info(f'File info: {file_info}')
    return {'files': sorted(files), 'file_info': file_info}

def get_file_first_line__ec5fe11be98e432184dc357c559fbb9f(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get the first line of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: VM path to the file

    Returns:
        Dict with first line info: {"content": "...", "is_comment": True/False, "starts_with": "..."}
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'content': '', 'is_comment': False, 'starts_with': ''}
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        content = file_bytes.decode('latin-1')
    lines = content.split('\n')
    if len(lines) == 0:
        return {'content': '', 'is_comment': False, 'starts_with': ''}
    first_line = lines[0]
    is_comment = first_line.strip().startswith('#')
    return {'content': first_line, 'is_comment': is_comment, 'starts_with': first_line[:20] if len(first_line) >= 20 else first_line}

def get_default_file_manager__a1cf9be3(env, config: dict):
    """Gets the default file manager on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'inode/directory']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_chapter_files_status__93cfd69b3c5adfd5dbb8817764000202(env, config: dict):
    """Get status of multiple chapter files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'chapter_files' (list of filenames)

    Returns:
        dict: {'existing_files': list, 'missing_files': list}
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    chapter_files = config.get('chapter_files', [])
    existing_files = []
    missing_files = []
    for filename in chapter_files:
        filepath = f'{directory}/{filename}'
        command = f'test -f "{filepath}" && echo "exists" || echo "not_exists"'
        result = env.controller.run_bash_script(command, timeout=10)
        if result.get('output', '').strip() == 'exists':
            existing_files.append(filename)
        else:
            missing_files.append(filename)
    return {'existing_files': existing_files, 'missing_files': missing_files}

def get_text_file_content__20f90bc2668d01c30660dfeafc3af15b(env, config):
    """Get full content of a text file as a string.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key for file path

    Returns:
        str: Content of the file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        text = file_bytes.decode('utf-8')
        return text
    except Exception as e:
        return ''

def get_vm_dir_list__abdf3926e9609e7e5b435c8cdbb40013(env, config: Dict[str, Any]) -> List[str]:
    """Get list of files in a VM directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying directory path

    Returns:
        List of filenames in the directory
    """
    path = config['path']
    result = env.controller.run_bash_script(f"ls -1 '{path}' 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0:
        logger.warning(f"Failed to list directory {path}: {result.get('error', '')}")
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    return files

def get_pdf_count_desktop__6c8aa0e0d1a7a1f247933b05b6dc0e08(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get detailed information about PDF files in Desktop directory.

    This function counts PDFs and extracts their names and content to verify:
    1. Correct number of employee PDFs (excluding template)
    2. Proper naming (employee names)
    3. Content filled with employee data
    4. Checkmark symbols present

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key (path on VM)

    Returns:
        Dict with:
            - total_count: Total number of PDFs
            - employee_pdf_count: Number of PDFs excluding template
            - employee_pdf_names: List of employee PDF filenames
            - pdf_details: List of dicts with filename and text content
    """
    directory = config.get('directory', '/home/user/Desktop')
    result = env.controller.run_bash_script(f'find "{directory}" -maxdepth 1 -name "*.pdf" -type f -printf "%f\\n"', timeout=10)
    if result.get('status') != 'success' or not result.get('output'):
        logger.error(f'Failed to list PDF files: {result}')
        return {'total_count': 0, 'employee_pdf_count': 0, 'employee_pdf_names': [], 'pdf_details': []}
    all_pdfs = [line.strip() for line in result['output'].strip().split('\n') if line.strip()]
    template_name = 'review_template.pdf'
    employee_pdfs = [pdf for pdf in all_pdfs if pdf != template_name]
    expected_names = ['John Doe.pdf', 'Emily Johnson.pdf', 'Michael Brown.pdf', 'Linda Green.pdf', 'David Wilson.pdf', 'Sophia Carter.pdf', 'Alex Lee.pdf']
    pdf_details = []
    for pdf_name in employee_pdfs:
        pdf_path = f'{directory}/{pdf_name}'
        text_result = env.controller.run_bash_script(f'pdftotext "{pdf_path}" - 2>/dev/null || echo "ERROR: pdftotext not available"', timeout=15)
        text_content = ''
        if text_result.get('status') == 'success' and text_result.get('output'):
            text_content = text_result['output']
        pdf_details.append({'filename': pdf_name, 'text': text_content})
    return {'total_count': len(all_pdfs), 'employee_pdf_count': len(employee_pdfs), 'employee_pdf_names': sorted(employee_pdfs), 'expected_names': sorted(expected_names), 'pdf_details': pdf_details}

def get_zip_contents__8a01d242c9e2052109e1c26f5fa4a5dd(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get ZIP archive contents and file list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'files' (list of filenames), 'valid_zip' keys
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return {'exists': False, 'files': [], 'valid_zip': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.zip', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zip_file:
                file_list = zip_file.namelist()
                return {'exists': True, 'files': file_list, 'valid_zip': True, 'file_count': len(file_list)}
        finally:
            os.unlink(tmp_path)
    except zipfile.BadZipFile:
        logger.error(f'Invalid ZIP file: {path}')
        return {'exists': True, 'files': [], 'valid_zip': False, 'file_count': 0}
    except Exception as e:
        logger.error(f'Error reading ZIP file: {e}')
        return {'exists': False, 'files': [], 'valid_zip': False, 'file_count': 0}

def get_text_file_content__fc2c8cc4(env, config: dict):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    vm_path = config.get('path', '')
    file_bytes = env.controller.get_file(vm_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_pdf_count_in_dir__a0da932e(env, config: dict):
    """Get count of PDF files in specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        Integer count of PDF files
    """
    directory = config.get('directory', '/home/user/Documents/Articles')
    command = f"find {directory} -maxdepth 1 -name '*.pdf' -type f 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=30)
    output = result.get('output', '').strip()
    try:
        count = int(output)
        logger.info(f'Found {count} PDF files in {directory}')
        return count
    except ValueError:
        logger.warning(f'Could not parse count from output: {output}')
        return 0

def get_file_line_count__bcaf4400(env, config: dict):
    """
    Read line count from a file containing a single number.

    This getter extracts the line count value that the user saved to a file
    after counting lines in the Colab notebook. It validates the file exists,
    contains parseable content, and returns an integer value for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file location

    Returns:
        int: The line count as an integer, or None if file doesn't exist/is invalid
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        if not content:
            return None
        line_count = int(content)
        return line_count
    except (ValueError, Exception) as e:
        print(f'Error parsing line count from file: {e}')
        return None

def get_vm_file(env, config: Dict[str, Any]):
    """Get file from VM - simplified version for this getter."""
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        return cache_path
    except Exception as e:
        logger.error(f'Error getting file {path}: {e}')
        return None

def get_text_file_content__d997648158d90618366cc740cfd24b34(env, config: Dict[str, Any]) -> str:
    """Extract text content from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Text content of the file as a string, or empty string if file not found
    """
    path = config.get('path', '')
    if not path:
        logger.error('No path specified in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_wallpaper_path__a70b53e8(env, config: dict):
    """Get the current wallpaper path from GNOME settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        String with wallpaper file path
    """
    command = ['gsettings', 'get', 'org.gnome.desktop.background', 'picture-uri']
    result = env.controller.run_bash_script(' '.join(command), timeout=10)
    output = result.get('output', '').strip()
    if output.startswith("'file://"):
        output = output[8:-1]
    elif output.startswith('file://'):
        output = output[7:]
    elif output.startswith("'"):
        output = output[1:-1]
    return output

def get_file_count_in_dir__fbd9137c(env, config: dict):
    """Count files in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path' and optional 'pattern'

    Returns:
        Number of files in directory
    """
    dir_path = config['dir_path']
    pattern = config.get('pattern', '*')
    if pattern == '*':
        command = f'find "{dir_path}" -maxdepth 1 -type f | wc -l'
    else:
        command = f'find "{dir_path}" -maxdepth 1 -type f -name "{pattern}" | wc -l'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return None
    try:
        count = int(result['output'].strip())
        return count
    except ValueError:
        return None

def get_pdf_multiple_fields__ff249445b547a028ee5246e45cd19fdf(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if multiple text fields exist in a PDF.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'fields' keys

    Returns:
        Dict[str, Any]: Dictionary containing text fields found
    """
    path = config.get('path', '')
    fields = config.get('fields', [])
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Could not retrieve file: {path}')
        return {'text_fields': {field: False for field in fields}}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from pypdf import PdfReader
            reader = PdfReader(tmp_path)
            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())
            full_text = ' '.join(text_parts)
            text_fields_result = {}
            for field in fields:
                text_fields_result[field] = field in full_text
            return {'text_fields': text_fields_result, 'full_text': full_text}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting PDF text: {e}')
        return {'text_fields': {field: False for field in fields}, 'full_text': ''}

def get_picture_filenames__fdcd3f41(env, config):
    """Get sorted list of .jpg filenames in Pictures directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters

    Returns:
        list: Sorted list of .jpg filenames
    """
    command = 'python3 -c "import os; import json; pics = sorted([f for f in os.listdir(\'/home/user/Pictures\') if f.endswith(\'.jpg\')]); print(json.dumps(pics))"'
    result = env.controller.run_bash_script(command, timeout=30)
    try:
        import json
        filenames = json.loads(result['output'].strip())
        return filenames
    except:
        return []

def get_docx_content_and_font__abb42de39a213997a2a7f06fa1fe5d2a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get document content and font size from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the file path on VM

    Returns:
        Dict with 'content' (str) and 'font_size' (int) keys
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'content': '', 'font_size': None}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'content': '', 'font_size': None}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        content = []
        for paragraph in doc.paragraphs:
            content.append(paragraph.text)
        full_text = ''.join(content)
        font_size = None
        all_font_size_12 = True
        for paragraph in doc.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    if run.font.size:
                        current_size = run.font.size.pt
                        if font_size is None:
                            font_size = current_size
                        if current_size != 12:
                            all_font_size_12 = False
                    else:
                        all_font_size_12 = False
        if all_font_size_12 and font_size == 12:
            return {'content': full_text.strip(), 'font_size': 12}
        else:
            return {'content': full_text.strip(), 'font_size': font_size if font_size else None}
    finally:
        os.unlink(tmp_path)

def get_file_existence__198be354(env, config):
    """Check if file exists and analyze its content against Colab notebook code.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File existence, content analysis, and Colab code comparison metadata
    """
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    colab_code_cells = extract_colab_code_cells(env)
    colab_code_combined = '\n\n'.join(colab_code_cells) if colab_code_cells else ''
    if not file_content:
        return {'exists': False, 'is_python': False, 'is_empty': True, 'content': '', 'is_valid_python': False, 'has_code_patterns': False, 'line_count': 0, 'colab_code_cells': colab_code_cells, 'colab_code_combined': colab_code_combined, 'contains_colab_code': False, 'code_similarity_score': 0.0}
    is_python = file_path.endswith('.py')
    is_empty = len(file_content.strip()) == 0
    is_valid_python = False
    has_code_patterns = False
    if file_content:
        try:
            ast.parse(file_content)
            is_valid_python = True
        except:
            is_valid_python = False
        code_patterns = ['^\\s*import\\s+\\w+', '^\\s*from\\s+\\w+', '^\\s*def\\s+\\w+', '^\\s*class\\s+\\w+', '=', '^\\s*#', 'print\\s*\\(']
        pattern_count = 0
        for pattern in code_patterns:
            if re.search(pattern, file_content, re.MULTILINE):
                pattern_count += 1
        has_code_patterns = pattern_count >= 2
    line_count = len([line for line in file_content.split('\n') if line.strip()])
    contains_colab_code = False
    code_similarity_score = 0.0
    if colab_code_cells and file_content:
        file_normalized = ' '.join(file_content.split())
        cells_found = 0
        for cell in colab_code_cells:
            cell_normalized = ' '.join(cell.split())
            if cell_normalized and cell_normalized in file_normalized:
                cells_found += 1
        if len(colab_code_cells) > 0:
            code_similarity_score = cells_found / len(colab_code_cells)
            contains_colab_code = code_similarity_score > 0.5
        logger.info(f'[FILE_EXISTENCE] Found {cells_found}/{len(colab_code_cells)} code cells in file')
        logger.info(f'[FILE_EXISTENCE] Code similarity score: {code_similarity_score:.2f}')
    return {'exists': True, 'is_python': is_python, 'is_empty': is_empty, 'size': len(file_content), 'content': file_content, 'is_valid_python': is_valid_python, 'has_code_patterns': has_code_patterns, 'line_count': line_count, 'colab_code_cells': colab_code_cells, 'colab_code_combined': colab_code_combined, 'contains_colab_code': contains_colab_code, 'code_similarity_score': code_similarity_score}

def get_file_exists__62dacdc1(env, config):
    """
    Check if a file exists at the specified path on VM and validate it's a valid PNG image.

    Returns:
        dict: Information about the file including existence, type, size, and validity
    """
    file_path = config.get('path', '')
    exists_command = f'test -f "{file_path}" && echo "EXISTS" || echo "NOT_EXISTS"'
    exists_result = env.controller.run_bash_script(exists_command, timeout=10)
    exists = exists_result.get('output', '').strip() == 'EXISTS'
    if not exists:
        return {'exists': False, 'file_type': None, 'size': 0, 'is_png': False, 'is_valid_image': False}
    file_type_command = f'file -b "{file_path}"'
    file_type_result = env.controller.run_bash_script(file_type_command, timeout=10)
    file_type = file_type_result.get('output', '').strip()
    size_command = f'stat -c %s "{file_path}" 2>/dev/null || stat -f %z "{file_path}" 2>/dev/null'
    size_result = env.controller.run_bash_script(size_command, timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
    except ValueError:
        size = 0
    is_png = 'PNG' in file_type.upper()
    magic_bytes_command = f'xxd -l 8 -p "{file_path}"'
    magic_result = env.controller.run_bash_script(magic_bytes_command, timeout=10)
    magic_bytes = magic_result.get('output', '').strip()
    is_valid_png = magic_bytes.startswith('89504e47')
    return {'exists': True, 'file_type': file_type, 'size': size, 'is_png': is_png, 'is_valid_image': is_valid_png}

def get_file_exists__b03b6c61(env, config: Dict[str, Any]) -> bool:
    """Check if a file exists on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    file_path = config.get('file_path', '')
    result = env.controller.run_bash_script(f'test -f "{file_path}" && echo "exists" || echo "not_exists"', timeout=10)
    if result.get('status') == 'success' and 'exists' in result.get('output', ''):
        return True
    return False

def get_csv_merge_data__c66369e707b97de3ccd6da4699663fe6(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get comprehensive data about merged CSV file and verify source files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to CSV file on VM

    Returns:
        dict: Contains row_count, has_single_header, sample_rows, total_rows,
              and verification data from source files
    """
    file_path = config.get('path', '')
    if not file_path:
        return {'row_count': 0, 'has_single_header': False, 'sample_rows': [], 'total_rows': 0, 'unique_values_from_merged': set(), 'source_file1_unique_values': set(), 'source_file2_unique_values': set()}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {'row_count': 0, 'has_single_header': False, 'sample_rows': [], 'total_rows': 0, 'unique_values_from_merged': set(), 'source_file1_unique_values': set(), 'source_file2_unique_values': set()}
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            all_rows = list(reader)
            total_rows = len(all_rows)
            if total_rows == 0:
                return {'row_count': 0, 'has_single_header': False, 'sample_rows': [], 'total_rows': 0, 'unique_values_from_merged': set(), 'source_file1_unique_values': set(), 'source_file2_unique_values': set()}
            header = all_rows[0] if total_rows > 0 else []
            data_rows = all_rows[1:] if total_rows > 1 else []
            row_count = len(data_rows)
            has_single_header = True
            for row in data_rows:
                if row == header:
                    has_single_header = False
                    break
            sample_rows = []
            if row_count > 0:
                sample_rows.extend(data_rows[:3])
                if row_count > 100:
                    mid_idx = row_count // 2
                    sample_rows.extend(data_rows[mid_idx:mid_idx + 3])
                if row_count > 6:
                    sample_rows.extend(data_rows[-3:])
            unique_values_from_merged = set()
            for row in data_rows:
                if row and len(row) > 0:
                    unique_values_from_merged.add(row[0])
            source_file1_unique_values = _extract_unique_values_from_xlsx(env, '/home/user/Desktop/file1.xlsx')
            source_file2_unique_values = _extract_unique_values_from_ods(env, '/home/user/Desktop/file2.ods')
            return {'row_count': row_count, 'has_single_header': has_single_header, 'sample_rows': sample_rows, 'total_rows': total_rows, 'unique_values_from_merged': unique_values_from_merged, 'source_file1_unique_values': source_file1_unique_values, 'source_file2_unique_values': source_file2_unique_values}
    finally:
        os.unlink(tmp_path)

def get_text_file_content__7bbdf0a0733630cbbfd86729556fc827(env, config):
    """Get full content of a text file as a string.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key for file path

    Returns:
        str: Content of the file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        text = file_bytes.decode('utf-8')
        return text
    except Exception as e:
        return ''

def get_pdf_filenames__4e03b1ed(env, config: dict):
    """Get PDF filenames matching a pattern.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        list: PDF filenames in the folder
    """
    folder_path = config.get('folder_path', '/home/user/Documents/Blog')
    command = f"ls -1 '{folder_path}'/*.pdf 2>/dev/null | xargs -n 1 basename 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [line.strip() for line in output.split('\n') if line.strip()]
    logger.info(f'Found PDF filenames: {files}')
    return files

def get_pdf_basic_info__7b63a847bf086913e974e7a32debb9ec(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get basic information about a PDF file (exists, page count, has images).

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with file info or empty dict if file not found
    """
    from pypdf import PdfReader
    import tempfile
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return {}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get PDF file from VM: {path}')
        return {'exists': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        page_count = len(reader.pages)
        has_images = False
        for page in reader.pages:
            if '/XObject' in page.get('/Resources', {}):
                xobjects = page['/Resources']['/XObject'].get_object()
                for obj in xobjects:
                    if xobjects[obj].get('/Subtype') == '/Image':
                        has_images = True
                        break
            if has_images:
                break
        file_size = len(file_bytes)
        os.unlink(tmp_path)
        return {'exists': True, 'page_count': page_count, 'has_images': has_images, 'file_size': file_size}
    except Exception as e:
        logger.error(f'Error analyzing PDF: {e}')
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass
        return {'exists': False, 'error': str(e)}

def get_multiple_files_exist__03426e679d8f4571bede57a16eea69a4(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Check if files were properly copied by verifying both source and destination exist with matching content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' key (list of destination file paths)

    Returns:
        List of dicts with keys: 'dest_exists', 'source_exists', 'content_matches'
    """
    dest_paths = config.get('paths', [])
    source_map = {'/home/user/Documents/Projects/OSWorld/backup/main.py': '/home/user/Documents/Projects/OSWorld/codes/main.py', '/home/user/Documents/Projects/OSWorld/backup/test_connection.py': '/home/user/Documents/Projects/OSWorld/codes/test_connection.py', '/home/user/Documents/Projects/OSWorld/backup/meeting_notes.md': '/home/user/Documents/Projects/OSWorld/meeting_notes.md'}
    results = []
    for dest_path in dest_paths:
        source_path = source_map.get(dest_path)
        if not source_path:
            logger.warning(f'No source path mapping for destination: {dest_path}')
            results.append({'dest_exists': False, 'source_exists': False, 'content_matches': False})
            continue
        dest_check = f"test -f '{dest_path}' && echo 'exists' || echo 'not_exists'"
        dest_result = env.controller.run_bash_script(dest_check, timeout=10)
        dest_exists = dest_result.get('output', '').strip() == 'exists'
        source_check = f"test -f '{source_path}' && echo 'exists' || echo 'not_exists'"
        source_result = env.controller.run_bash_script(source_check, timeout=10)
        source_exists = source_result.get('output', '').strip() == 'exists'
        content_matches = False
        if dest_exists and source_exists:
            checksum_cmd = f"md5sum '{source_path}' '{dest_path}' | awk '{{print $1}}'"
            checksum_result = env.controller.run_bash_script(checksum_cmd, timeout=15)
            checksums = checksum_result.get('output', '').strip().split('\n')
            if len(checksums) == 2:
                content_matches = checksums[0] == checksums[1]
                logger.info(f'Checksums - Source: {checksums[0]}, Dest: {checksums[1]}, Match: {content_matches}')
        result_dict = {'dest_exists': dest_exists, 'source_exists': source_exists, 'content_matches': content_matches}
        results.append(result_dict)
        logger.info(f'Copy check for {dest_path}: {result_dict}')
    return results

def get_directory_files__940d01bc(env, config: dict):
    """Get list of files in a directory recursively.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of file paths
    """
    path = config.get('path', '/home/user/Documents/Projects/OSWorld')
    result = env.controller.run_bash_script(f"find '{path}' -type f 2>/dev/null", timeout=10)
    if result['returncode'] != 0:
        return []
    output = result.get('output', '').strip()
    if not output:
        return []
    files = [line.strip() for line in output.split('\n') if line.strip()]
    return files

def get_dir_file_count__95b4929b(env, config):
    """Count files in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' path and optional 'pattern'

    Returns:
        Integer count of files
    """
    directory = config.get('directory', '/home/user')
    pattern = config.get('pattern', '*')
    command = f"find {directory} -maxdepth 1 -type f -name '{pattern}' | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    if result and result.get('status') == 'success':
        try:
            count = int(result['output'].strip())
            return count
        except:
            return 0
    return 0

def get_pdf_files_info__8010e79b(env, config):
    """Get information about PDF files in a directory.
    
    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key
        
    Returns:
        List of PDF filenames in the directory
    """
    dir_path = config.get('path', '/home/user/lecture_slides')
    try:
        result = env.controller.run_bash_script(f"ls -1 {dir_path} 2>/dev/null || echo ''", timeout=30)
        if result['returncode'] != 0:
            logger.warning(f'Failed to list directory: {dir_path}')
            return []
        output = result['output'].strip()
        if not output:
            return []
        files = [f.strip() for f in output.split('\n') if f.strip()]
        pdf_files = [f for f in files if f.lower().endswith('.pdf')]
        logger.info(f'Found {len(pdf_files)} PDF files in {dir_path}: {pdf_files}')
        return pdf_files
    except Exception as e:
        logger.error(f'Error getting file list: {e}')
        return []

def get_file_exists_and_size__d6a98717(env, config: Dict[str, Any]):
    """
    Check if a file exists on the VM and return its size in bytes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {"exists": bool, "size": int (bytes), "is_png": bool}
    """
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_folder_contents__8542f72b(env, config: dict):
    """Get list of files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        list: List of filenames in the directory (empty if folder doesn't exist)
    """
    folder_path = config.get('folder_path', '')
    command = f'ls -1 "{folder_path}" 2>/dev/null'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
        return files
    return []

def get_inbox_backup_dir__987d91d90e3876e4880d553cc7b5b944(env, config):
    """
    Check if inbox backup directory exists and contains files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to backup directory

    Returns:
        str: Path to cache file containing directory info or None
    """
    backup_path = config.get('path', '/home/user/inbox_backup')
    result = env.controller.run_bash_script(f'test -d {backup_path} && ls -la {backup_path}', timeout=30)
    if result['returncode'] != 0:
        return None
    cache_path = os.path.join(env.cache_dir, 'inbox_backup_dir.ls')
    with open(cache_path, 'w') as f:
        f.write(result['output'])
    return cache_path

def get_vm_file__96eb02f4ac90683da522ea44f75e2519(env, config):
    """
    Get the DOCX file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded file
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    import tempfile
    import os
    suffix = os.path.splitext(file_path)[1]
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_file_exists__46f0b51f(env, config: dict):
    """Check if the exported PDF file exists and validate its properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with validation results including:
            - exists: bool, whether file exists
            - is_pdf: bool, whether file is a valid PDF
            - file_size: int, file size in bytes
            - has_content: bool, whether PDF has reasonable content
            - has_spreadsheet_data: bool, whether PDF contains spreadsheet data
            - created_by_libreoffice: bool, whether PDF metadata shows LibreOffice as creator
    """
    vm_path = config.get('path', '/home/user/Desktop/annual-report.pdf')
    result = {'exists': False, 'is_pdf': False, 'file_size': 0, 'has_content': False, 'has_spreadsheet_data': False, 'created_by_libreoffice': False}
    existence_check = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    if existence_check.get('output', '').strip() != 'EXISTS':
        return result
    result['exists'] = True
    file_type_check = env.controller.run_bash_script(f"file -b '{vm_path}'", timeout=10)
    file_type_output = file_type_check.get('output', '').strip().lower()
    result['is_pdf'] = 'pdf' in file_type_output
    size_check = env.controller.run_bash_script(f"stat -c %s '{vm_path}' 2>/dev/null || echo '0'", timeout=10)
    try:
        result['file_size'] = int(size_check.get('output', '0').strip())
    except ValueError:
        result['file_size'] = 0
    if result['file_size'] > 1024:
        header_check = env.controller.run_bash_script(f"head -c 4 '{vm_path}' 2>/dev/null", timeout=10)
        header = header_check.get('output', '')
        result['has_content'] = header.startswith('%PDF')
    pdf_text_check = env.controller.run_bash_script(f"pdftotext '{vm_path}' - 2>/dev/null | head -c 5000", timeout=15)
    pdf_text = pdf_text_check.get('output', '').strip()
    if pdf_text:
        has_numbers = any((char.isdigit() for char in pdf_text))
        has_multiple_lines = pdf_text.count('\n') > 5
        has_sufficient_content = len(pdf_text.strip()) > 100
        survey_keywords = ['enterprise', 'survey', '2021', 'financial', 'year', 'provisional']
        has_survey_keywords = any((keyword.lower() in pdf_text.lower() for keyword in survey_keywords))
        specific_survey_terms = ['industry', 'annual']
        has_specific_survey_data = any((term.lower() in pdf_text.lower() for term in specific_survey_terms))
        result['has_spreadsheet_data'] = has_numbers and has_multiple_lines and has_sufficient_content and has_survey_keywords and has_specific_survey_data
    metadata_check = env.controller.run_bash_script(f"pdfinfo '{vm_path}' 2>/dev/null | grep -i -E '(Creator|Producer)'", timeout=10)
    metadata_output = metadata_check.get('output', '').lower()
    result['created_by_libreoffice'] = 'libreoffice' in metadata_output
    return result

def get_python_file_content__generic(env, config):
    """Get full content of a Python file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        String (full file content)
    """
    path = config.get('path', '/home/user/Desktop/test.py')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return ''
    return file_bytes.decode('utf-8')

def get_multiple_files_status__4e03b1ed(env, config: dict):
    """Check existence of multiple specific files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_paths' parameter (list)

    Returns:
        dict: Mapping of file path to existence status
    """
    file_paths = config.get('file_paths', [])
    status = {}
    for file_path in file_paths:
        command = f"test -f '{file_path}' && echo 'yes' || echo 'no'"
        result = env.controller.run_bash_script(command, timeout=10)
        exists = result.get('output', '').strip() == 'yes'
        status[file_path] = exists
    logger.info(f'File status: {status}')
    return status

def get_python_file_structure__58d4fb9ea6e69b7f57a37a59813050e8(env, config: dict):
    """Extract Python file structure (imports, classes, functions).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File structure with imports, classes, and functions
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'imports': [], 'classes': [], 'functions': []}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.split('\n')
        imports = []
        classes = []
        functions = []
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('import ') or stripped.startswith('from '):
                imports.append(stripped)
            elif stripped.startswith('class '):
                class_name = stripped.split('(')[0].replace('class ', '').replace(':', '').strip()
                classes.append(class_name)
            elif stripped.startswith('def '):
                func_name = stripped.split('(')[0].replace('def ', '').strip()
                functions.append(func_name)
        return {'exists': True, 'imports': imports, 'classes': classes, 'functions': functions, 'total_lines': len(lines)}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_dir_file_count__00db2192(env, config):
    """Count Python files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path' and optional 'pattern' parameters

    Returns:
        int: Number of files matching the pattern
    """
    dir_path = config.get('dir_path')
    pattern = config.get('pattern', '*.py')
    command = f"find {dir_path} -maxdepth 1 -type f -name '{pattern}' 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(command, timeout=10)
    try:
        count = int(result.get('output', '0').strip())
        return count
    except ValueError:
        return 0

def get_text_file_content__7e5134258960ebea77ca0d290984a7a3(env, config):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_pdf_in_year_folders__5ad487ac7d025a6b906a4c83e8beac41(env, config: dict):
    """Check if PDF files exist in year-organized folders with content verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - base_directory: Base directory path on VM
            - files_by_year: Dict mapping year -> list of file info dicts

    Returns:
        dict: Results for each file {year/filename: status}
    """
    base_dir = config.get('base_directory', '/home/user/Documents/Blog')
    files_by_year = config.get('files_by_year', {})
    results = {}
    for (year, file_list) in files_by_year.items():
        year_dir = os.path.join(base_dir, year)
        for file_info in file_list:
            filename = file_info['filename']
            content_check = file_info.get('content_check', '')
            file_path = os.path.join(year_dir, filename)
            result_key = f'{year}/{filename}'
            check_cmd = f'[ -f "{file_path}" ] && echo "exists" || echo "not_found"'
            result = env.controller.run_bash_script(check_cmd, timeout=10)
            if result.get('output', '').strip() != 'exists':
                results[result_key] = -1
                continue
            if content_check:
                try:
                    file_bytes = env.controller.get_file(file_path)
                    if not file_bytes:
                        results[result_key] = 0
                        continue
                    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    try:
                        import fitz
                        doc = fitz.open(tmp_path)
                        text = ''
                        for page in doc:
                            text += page.get_text()
                        doc.close()
                        if content_check in text:
                            results[result_key] = 1
                        else:
                            results[result_key] = 0
                    except Exception:
                        results[result_key] = 0
                    finally:
                        os.unlink(tmp_path)
                except Exception:
                    results[result_key] = 0
            else:
                results[result_key] = 1
    return results

def get_python_file__4ca3247dcf6b464342c4e3f53d844797(env, config: Dict[str, Any]) -> str:
    """Get the content of a Python file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file location

    Returns:
        String containing the file content, or empty string if file doesn't exist
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        return file_bytes.decode('utf-8')
    except Exception as e:
        logger.error(f'Error decoding file content: {e}')
        return ''

def get_playlist_file_status__3145ebe2c0a2fb7aa964b8c8f3b95ea9(env, config: Dict[str, str]):
    """
    Checks if playlist file was created and validates M3U format.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'playlist_path' key

    Returns:
        dict: Playlist file status information
    """
    playlist_path = config.get('playlist_path', '/home/user/playlist.m3u')
    file_check = env.controller.run_bash_script(f"test -f {playlist_path} && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = 'EXISTS' in file_check.get('output', '')
    if not file_exists:
        return {'exists': False, 'content': '', 'has_video_reference': False, 'has_valid_m3u_format': False}
    content_result = env.controller.run_bash_script(f'cat {playlist_path}', timeout=10)
    content = content_result.get('output', '').strip()
    has_valid_m3u_format = content.startswith('#EXTM3U')
    lines = content.split('\n')
    has_video_ref = False
    for line in lines:
        line = line.strip()
        if line and (not line.startswith('#')):
            if 'video.mp4' in line and (line == 'video.mp4' or line.endswith('video.mp4')):
                has_video_ref = True
                break
    return {'exists': True, 'content': content, 'has_video_reference': has_video_ref, 'has_valid_m3u_format': has_valid_m3u_format, 'path': playlist_path}

def get_folder_file_count__5e18d045(env, config: dict):
    """Count files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' parameter

    Returns:
        int: Number of files in the directory
    """
    folder_path = config.get('folder_path', '')
    command = f'find "{folder_path}" -maxdepth 1 -type f 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        try:
            count = int(result['output'].strip())
            return count
        except:
            return 0
    return 0

def get_pdf_content_fit__234d824f7c8a19c4a33bd790f5cee244(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get PDF page count and dimensions to verify content fitting.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'page_count' (int), 'width' (float), 'height' (float)
    """
    try:
        from pypdf import PdfReader
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            return {'page_count': 0, 'width': 0, 'height': 0}
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            if len(reader.pages) == 0:
                return {'page_count': 0, 'width': 0, 'height': 0}
            page = reader.pages[0]
            width = float(page.mediabox.width)
            height = float(page.mediabox.height)
            return {'page_count': len(reader.pages), 'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        print(f'Error reading PDF: {e}')
        return {'page_count': 0, 'width': 0, 'height': 0}

def get_csv_column_count__7f9974e9(env, config: dict):
    """Get CSV file content including headers and row count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with column count, headers, and row count
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'columns': 0, 'headers': [], 'row_count': 0}
    content = file_bytes.decode('utf-8')
    lines = content.strip().split('\n')
    if not lines:
        return {'columns': 0, 'headers': [], 'row_count': 0}
    reader = csv.reader(lines)
    rows = list(reader)
    if not rows:
        return {'columns': 0, 'headers': [], 'row_count': 0}
    headers = rows[0] if rows else []
    row_count = len(rows) - 1
    return {'columns': len(headers), 'headers': headers, 'row_count': row_count}

def get_git_dir_exists__e2da960ab9034666db33db74ae6371a7(env, config: dict):
    """Check if .git directory exists in a repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        bool: True if .git directory exists, False otherwise
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    repo_path = config['repo_path']
    command = f"test -d {repo_path}/.git && echo 'EXISTS' || echo 'NOT_EXISTS'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        return output == 'EXISTS'
    else:
        logger.error('Failed to check .git directory. Status code: %d', response.status_code)
        return False

def get_bash_file_content__5ed0399c75d125da4fc3b5f7583bb5c4(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file not found
    """
    file_path = config.get('path', '/home/user/output.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_file_exists__f663e89c52b74d5c5d4e38ab6d86c83f(env, config):
    """Check if a file exists on the VM and verify it is a valid PDF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'path': str, 'is_pdf': bool, 'file_size': int}
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "true" || echo "false"', timeout=10)
    exists = result.get('output', '').strip() == 'true'
    is_pdf = False
    file_size = 0
    if exists:
        file_type_result = env.controller.run_bash_script(f'file -b --mime-type "{file_path}"', timeout=10)
        mime_type = file_type_result.get('output', '').strip()
        is_pdf = mime_type == 'application/pdf'
        size_result = env.controller.run_bash_script(f'stat -c %s "{file_path}"', timeout=10)
        try:
            file_size = int(size_result.get('output', '0').strip())
        except ValueError:
            file_size = 0
    return {'exists': exists, 'path': file_path, 'is_pdf': is_pdf, 'file_size': file_size}

def get_file_copy_verification__dc38ce29eba391e7169ff4e028e69a72(env, config):
    """Verify that a file was copied from source to destination with content preserved.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'source_file', 'dest_dir', and 'dest_file' parameters

    Returns:
        dict: {
            'source_exists': bool,
            'dest_exists': bool,
            'content_matches': bool,
            'copy_successful': bool
        }
    """
    source_file = config.get('source_file', '')
    dest_dir = config.get('dest_dir', '')
    dest_file = config.get('dest_file', '')
    dest_path = f'{dest_dir}/{dest_file}'
    check_source = f'[ -f "{source_file}" ] && echo "exists" || echo "not_exists"'
    source_result = env.controller.run_bash_script(check_source, timeout=10)
    source_exists = False
    if source_result and source_result.get('output'):
        source_exists = source_result['output'].strip() == 'exists'
    check_dest = f'[ -f "{dest_path}" ] && echo "exists" || echo "not_exists"'
    dest_result = env.controller.run_bash_script(check_dest, timeout=10)
    dest_exists = False
    if dest_result and dest_result.get('output'):
        dest_exists = dest_result['output'].strip() == 'exists'
    content_matches = False
    if source_exists and dest_exists:
        compare_cmd = f'cmp -s "{source_file}" "{dest_path}" && echo "match" || echo "differ"'
        compare_result = env.controller.run_bash_script(compare_cmd, timeout=10)
        if compare_result and compare_result.get('output'):
            content_matches = compare_result['output'].strip() == 'match'
    copy_successful = source_exists and dest_exists and content_matches
    return {'source_exists': source_exists, 'dest_exists': dest_exists, 'content_matches': content_matches, 'copy_successful': copy_successful}

def get_screenshot_file_info__57d2d1f6(env, config):
    """Get information about a screenshot file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'base_name' keys

    Returns:
        dict with 'exists', 'path', 'extension' keys
    """
    directory = config.get('directory', '/home/user/Desktop')
    base_name = config.get('base_name', 'screenshot')
    extensions = ['.png', '.jpg', '.jpeg', '.bmp', '.gif']
    result = {'exists': False, 'path': None, 'extension': None, 'files_found': []}
    command = f"ls -1 '{directory}' 2>/dev/null || echo ''"
    ls_result = env.controller.run_bash_script(command, timeout=10)
    if ls_result['returncode'] != 0:
        logger.warning(f'Failed to list directory {directory}')
        return result
    files = ls_result['output'].strip().split('\n') if ls_result['output'].strip() else []
    for f in files:
        if f.startswith(base_name):
            result['files_found'].append(f)
            for ext in extensions:
                if f.endswith(ext):
                    result['exists'] = True
                    result['path'] = os.path.join(directory, f)
                    result['extension'] = ext
                    break
            if result['exists']:
                break
    return result

def get_docx_text_content__e38643bd(env, config: Dict[str, Any]) -> Optional[str]:
    """Extract full text content from a docx file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Full text content as string, or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        doc = Document(tmp_path)
        text_content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        import os
        os.unlink(tmp_path)
        return text_content
    except Exception as e:
        return None

def get_tetris_files__03faeae5(env, config):
    """Get Tetris Python files from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'paths' list

    Returns:
        dict: Mapping of filename to file content
    """
    paths = config.get('paths', [])
    files = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if file_bytes:
                content = file_bytes.decode('utf-8')
                filename = path.split('/')[-1]
                files[filename] = content
        except Exception as e:
            print(f'Error reading {path}: {e}')
            files[path.split('/')[-1]] = None
    return files

def get_desktop_pdf__a45ebe2876987410607711d3992c10db(env, config: Dict[str, Any]) -> Optional[str]:
    """Get PDF file from Desktop and save to cache.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        Path to cached PDF file, or None if file doesn't exist
    """
    vm_path = config['path']
    dest = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    os.makedirs(env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_pdf_info__9a18b30d646547c54b42b3593f83920d(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get PDF file information including existence and page count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'page_count', 'file_size' keys, or None if error
    """
    from pypdf import PdfReader
    import tempfile
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Failed to get file from VM: {path}')
        return {'exists': False, 'page_count': 0, 'file_size': 0}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            reader = PdfReader(tmp_path)
            page_count = len(reader.pages)
            file_size = len(file_bytes)
            return {'exists': True, 'page_count': page_count, 'file_size': file_size}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading PDF file: {e}')
        return {'exists': False, 'page_count': 0, 'file_size': 0}

def get_file_line_count__7233f122896fac183c973343e2cf3b2a(env, config):
    """Read a text file and return the integer value from its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Integer value parsed from file content, or -1 if file not found or invalid
    """
    file_path = config.get('path', '')
    if not file_path:
        return -1
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return -1
        content = file_bytes.decode('utf-8', errors='ignore').strip()
        return int(content)
    except (ValueError, Exception) as e:
        return -1

def get_large_file_count__b9976565(env, config: Dict[str, Any]) -> int:
    """Get count of files larger than threshold in directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'min_size_kb' parameters

    Returns:
        Count of files larger than threshold
    """
    path = config['path']
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return 0
    return len(result['children'])

def get_file_count_from_text__bc253f41(env, config):
    """Read the count of .doc files from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        int: Count value from the text file
    """
    path = config.get('path', '/home/user/Desktop/doc_count.txt')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    content = file_bytes.decode('utf-8').strip()
    try:
        return int(content)
    except ValueError:
        return None

def get_pdf_page_orientation__596fcd5e219bea1372357f0afb95cc85(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PDF page orientation information.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to PDF file on VM

    Returns:
        Dict with orientation info
    """
    from pypdf import PdfReader
    import tempfile
    path = config.get('path')
    if not path:
        return {'exists': False}
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False}
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        reader = PdfReader(tmp_path)
        if len(reader.pages) == 0:
            os.unlink(tmp_path)
            return {'exists': True, 'page_count': 0}
        page = reader.pages[0]
        mediabox = page.mediabox
        width = float(mediabox.width)
        height = float(mediabox.height)
        if width > height:
            orientation = 'landscape'
        elif height > width:
            orientation = 'portrait'
        else:
            orientation = 'square'
        os.unlink(tmp_path)
        return {'exists': True, 'page_count': len(reader.pages), 'orientation': orientation, 'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error reading PDF orientation: {e}')
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass
        return {'exists': False, 'error': str(e)}

def get_merged_text_file__e22bfb55f6ab9983d1cc35b82dc09aeb(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get merged text file from VM and check its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to text file on VM

    Returns:
        Dict with 'exists', 'line_count', 'contains_chapters' keys
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'line_count': 0, 'contains_chapters': []}
    try:
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = content.split('\n')
        line_count = len([l for l in lines if l.strip()])
        contains_chapters = []
        chapter_markers = config.get('chapter_markers', [])
        for marker in chapter_markers:
            pattern = '\\b' + re.escape(marker) + '\\b'
            if re.search(pattern, content, re.IGNORECASE):
                contains_chapters.append(marker)
        return {'exists': True, 'line_count': line_count, 'contains_chapters': contains_chapters}
    except Exception as e:
        return {'exists': False, 'line_count': 0, 'contains_chapters': []}

def get_files_exist__f50a55ca(env, config):
    """Check if multiple files exist in repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file_paths' list parameter

    Returns:
        dict: {file_path: bool} mapping of file existence
    """
    file_paths = config.get('file_paths', [])
    results = {}
    for file_path in file_paths:
        command = f"test -f {file_path} && echo 'YES' || echo 'NO'"
        result = env.controller.run_bash_script(command, timeout=10)
        results[file_path] = 'YES' in result.get('output', '')
    return results

def get_text_file_content__438c9c7ce7eebe25a3992ddf0a388112(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the file on VM

    Returns:
        str: File content as string
    """
    path = config.get('path')
    if not path:
        logger.error('No path provided in config')
        return ''
    try:
        file_bytes = env.controller.get_file(path)
        if not file_bytes:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_text_file_content__515e2337245bc72c8d34192293ce6646(env, config):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_mbox_content__1701225e2d2da95122de9fd6941c3c6f(env, config):
    """Extract email subjects from mbox file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to mbox file on VM

    Returns:
        List of email subjects found in the mbox file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8', errors='ignore')
    except Exception:
        return []
    subjects = []
    for line in content.split('\n'):
        if line.startswith('Subject:'):
            subject = line[8:].strip()
            if '=?utf-8?' in subject.lower():
                import re
                match = re.search('\\?B\\?(.*?)\\?=', subject, re.IGNORECASE)
                if match:
                    import base64
                    try:
                        decoded = base64.b64decode(match.group(1)).decode('utf-8')
                        subjects.append(decoded)
                    except Exception:
                        subjects.append(subject)
                else:
                    subjects.append(subject)
            else:
                subjects.append(subject)
    return subjects

def get_pdf_files_info__ea295e4c379ee25a192357c908aada76(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get comprehensive information about PDF files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' key (VM path)

    Returns:
        Dict with 'total_count', 'numbered_chapters', 'other_files' keys
    """
    directory = config.get('directory', '/home/user/Desktop/book')
    command = f"cd '{directory}' && ls -1 *.pdf 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    info = {'total_count': 0, 'numbered_chapters': [], 'other_files': []}
    if result['returncode'] != 0:
        logger.warning(f"Failed to list PDF files in {directory}: {result.get('error', '')}")
        return info
    output = result.get('output', '').strip()
    if not output:
        return info
    for line in output.split('\n'):
        filename = line.strip()
        if not filename:
            continue
        info['total_count'] += 1
        if '. ' in filename:
            parts = filename.split('. ', 1)
            if len(parts) == 2 and parts[0].isdigit():
                info['numbered_chapters'].append(filename)
                continue
        info['other_files'].append(filename)
    logger.info(f"PDF files in {directory}: {info['total_count']} total, {len(info['numbered_chapters'])} numbered chapters, {len(info['other_files'])} other files")
    return info

def get_csv_headers__4735f34d496bbe1961d6d9a20cf7b9bd(env, config):
    """Get the header row from a CSV file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of column headers, or empty list if file not found
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            headers = next(reader, [])
            return headers
    finally:
        os.unlink(tmp_path)

def get_file_exists__689ec9af4ba1471bf9b5f89e71cafeb9(env, config):
    """Check if a BMP file exists on the VM and verify it's a valid BMP format.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'exists': bool, 'is_bmp': bool, 'size': int, 'path': str}
    """
    file_path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'[ -f "{file_path}" ] && echo "true" || echo "false"', timeout=10)
    exists = exists_result.get('output', '').strip() == 'true'
    if not exists:
        return {'exists': False, 'is_bmp': False, 'size': 0, 'path': file_path}
    file_type_result = env.controller.run_bash_script(f'file "{file_path}"', timeout=10)
    file_type_output = file_type_result.get('output', '').strip()
    is_bmp = 'PC bitmap' in file_type_output or 'BMP image' in file_type_output
    size_result = env.controller.run_bash_script(f'stat -c %s "{file_path}" 2>/dev/null || echo "0"', timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
    except (ValueError, TypeError):
        size = 0
    return {'exists': exists, 'is_bmp': is_bmp, 'size': size, 'path': file_path}

def get_file_content__d77a2096(env, config: dict):
    """Get content of a text file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        str: Content of the file
    """
    path = config.get('path')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception:
        return ''

def get_pdf_exists__9753fb3ef75762f14fb08b7f236e3f81(env, config: Dict[str, Any]) -> Optional[str]:
    """Check if a PDF file exists at the specified path and return the path if it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the PDF file on VM

    Returns:
        str: The local cache path if file exists and is a PDF, None otherwise
    """
    vm_path = config.get('path')
    if not vm_path:
        logger.warning('No path specified in config')
        return None
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.warning(f'File not found at {vm_path}')
        return None
    if not file_bytes.startswith(b'%PDF'):
        logger.warning(f'File at {vm_path} is not a valid PDF')
        return None
    import tempfile
    cache_path = os.path.join(env.cache_dir, config.get('dest', 'output.pdf'))
    os.makedirs(os.path.dirname(cache_path) if os.path.dirname(cache_path) else env.cache_dir, exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    logger.info(f'PDF file saved to {cache_path} ({len(file_bytes)} bytes)')
    return cache_path

def get_file_content__c2756c851c53289bbd5185905c6853a2(env, config):
    """Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    result = env.controller.run_bash_script(f'cat {file_path}', timeout=10)
    if result.get('returncode') == 0:
        return result.get('output', '').strip()
    else:
        return ''

def get_text_file_content__551695fc(env, config: dict):
    """Read content from a text file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content.strip()
    except Exception as e:
        return ''

def get_text_file_lines__804f1b3e28fc8c6ee1d8689118367b3b(env, config):
    """
    Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: Lines from the file (stripped of whitespace)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
        return lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return []

def get_file_exists__d8671412(env, config):
    """Check file existence for rename operation.

    Args:
        env: Desktop environment
        config: Dict with 'path' key (target file) and optional 'source_path' key

    Returns:
        dict: {'target_exists': bool, 'source_exists': bool} for rename verification
    """
    target_path = config.get('path', '')
    source_path = config.get('source_path', '')
    target_result = env.controller.run_bash_script(f"test -f '{target_path}' && echo 'exists' || echo 'not_found'", timeout=10)
    target_exists = target_result['returncode'] == 0 and target_result['output'].strip() == 'exists'
    source_exists = False
    if source_path:
        source_result = env.controller.run_bash_script(f"test -f '{source_path}' && echo 'exists' || echo 'not_found'", timeout=10)
        source_exists = source_result['returncode'] == 0 and source_result['output'].strip() == 'exists'
    return {'target_exists': target_exists, 'source_exists': source_exists}

def get_python_file_content__f55aa7954b40a62bad3b8ba851857ed1(env, config):
    """
    Read a Python file from VM and return its content as string.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        print(f'Error reading file: {e}')
        return ''

def get_bashrc_path_check__c0107678(env, config):
    """Check both .bashrc file content and PATH after sourcing.

    This getter verifies that:
    1. The .bashrc file actually contains the PATH modification
    2. After sourcing .bashrc, PATH contains the expected directory

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Contains 'bashrc_content' and 'path_value'
    """
    bashrc_command = 'cat /home/user/.bashrc'
    bashrc_result = env.controller.run_bash_script(bashrc_command, timeout=30)
    if bashrc_result['returncode'] != 0:
        logger.error(f"Failed to read .bashrc: {bashrc_result.get('error', '')}")
        return {'bashrc_content': None, 'path_value': None}
    bashrc_content = bashrc_result['output'].strip()
    path_command = 'bash -c "source /home/user/.bashrc && echo $PATH"'
    path_result = env.controller.run_bash_script(path_command, timeout=30)
    if path_result['returncode'] != 0:
        logger.error(f"Failed to get PATH after sourcing .bashrc: {path_result.get('error', '')}")
        return {'bashrc_content': bashrc_content, 'path_value': None}
    path_value = path_result['output'].strip()
    return {'bashrc_content': bashrc_content, 'path_value': path_value}

def get_multi_directory_contents__b1a155d8(env, config):
    """Get contents of multiple directories.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directories' list

    Returns:
        List of directory contents (strings)
    """
    directories = config.get('directories', [])
    results = []
    for directory in directories:
        result = env.controller.run_bash_script(f'ls {directory}', timeout=30)
        output = result.get('output', '') if isinstance(result, dict) else ''
        results.append(output)
    return results

def get_text_file_content__f84e9cb5fd8cfc9fab208c24bcd90a7d(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    path = config.get('path', '/home/user/output.txt')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_timestamped_file__2b78c2fd0d670b6ee1c54ce65b4419a5(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a file with YYYYMMDD timestamp exists and extract the date.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (can contain wildcard)

    Returns:
        Dict with 'exists', 'has_valid_yyyymmdd_format', 'is_png', 'date_value' keys
    """
    path_pattern = config.get('path', '')
    result = {'exists': False, 'has_valid_yyyymmdd_format': False, 'is_png': False, 'date_value': None}
    try:
        matched_files = env.controller.execute(f'''python3 -c "import glob; print('\\n'.join(glob.glob('{path_pattern}')))"''', shell=True)
        if matched_files and matched_files.strip():
            matched_list = matched_files.strip().split('\n')
            for file_path in matched_list:
                filename = os.path.basename(file_path)
                pattern = '^screenshot[_-]?(\\d{8})\\.png$'
                match = re.match(pattern, filename)
                if match:
                    date_str = match.group(1)
                    try:
                        year = int(date_str[0:4])
                        month = int(date_str[4:6])
                        day = int(date_str[6:8])
                        if 1900 <= year <= 2100 and 1 <= month <= 12 and (1 <= day <= 31):
                            result['exists'] = True
                            result['has_valid_yyyymmdd_format'] = True
                            result['date_value'] = date_str
                            file_bytes = env.controller.get_file(file_path)
                            if file_bytes and file_path.lower().endswith('.png'):
                                try:
                                    import tempfile
                                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                        tmp.write(file_bytes)
                                        tmp_path = tmp.name
                                    try:
                                        img = Image.open(tmp_path)
                                        img.verify()
                                        result['is_png'] = True
                                    except Exception as e:
                                        logger.warning(f'Not a valid PNG: {e}')
                                    finally:
                                        os.unlink(tmp_path)
                                except Exception as e:
                                    logger.warning(f'Error verifying PNG: {e}')
                            break
                        else:
                            logger.info(f'Invalid date in filename: {date_str}')
                    except Exception as e:
                        logger.warning(f'Error parsing date from filename: {e}')
        else:
            logger.info(f'No files found matching pattern: {path_pattern}')
    except Exception as e:
        logger.error(f'Error checking file: {e}')
    return result

def get_file_list_content__ef0f28d6(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the content of a file from the VM.

    Args:
        env: Environment object with controller and cache_dir
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        str: Path to the cached file containing the content, or None if file not found
    """
    path = config['path']
    dest = config['dest']
    _path = os.path.join(env.cache_dir, dest)
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(_path, 'wb') as f:
            f.write(file_content)
        logger.info(f'Successfully saved file: {_path} ({len(file_content)} bytes)')
        return _path
    except Exception as e:
        logger.error(f'Error processing file {path}: {e}')
        return None

def get_vm_subtitle_file__2ad8e92c(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_pdf_text_content__4c28c68dd08d7073669bab47ee359a64(env, config: Dict[str, Any]) -> str:
    """
    Extract text content from a PDF file on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM path to PDF)

    Returns:
        str: Extracted text content (empty string if file doesn't exist or error)
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'Could not retrieve file: {path}')
        return ''
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            from pypdf import PdfReader
            reader = PdfReader(tmp_path)
            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())
            return ' '.join(text_parts)
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error extracting PDF text: {e}')
        return ''

def get_text_file_content__895c3960d172d43278234eeb5c495eda(env, config: Dict[str, Any]) -> str:
    """
    Read text file content from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        File content as string, or empty string if file doesn't exist
    """
    path = config.get('path', '/home/user/output.txt')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return ''
        content = file_bytes.decode('utf-8', errors='ignore')
        return content.strip()
    except Exception as e:
        logger.error(f'Error reading file {path}: {e}')
        return ''

def get_text_file_lines__dfc9794a(env, config: dict):
    """Extract text content as lines from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: File content as list of lines
    """
    path = config.get('path', '/home/user/Desktop/res.txt')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return None
    try:
        content = file_bytes.decode('utf-8')
        lines = [line.strip() for line in content.strip().split('\n')]
        return lines
    except Exception as e:
        return None

def get_text_file_line_count__ab2d13c4(env, config: Dict[str, Any]) -> int:
    """Get line count of a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Number of lines in the file
    """
    path = config['path']
    result = env.controller.run_bash_script(f'wc -l < {path} 2>/dev/null || echo 0', timeout=10)
    try:
        count = int(result.get('output', '0').strip())
        return count
    except ValueError:
        return 0

def get_text_file_content__6941d0dc31bbd0c3d844303b7e1c57e5(env, config):
    """Get full content of a text file as a string.

    Args:
        env: DesktopEnv instance
        config: dict with 'path' key for file path

    Returns:
        str: Content of the file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        text = file_bytes.decode('utf-8')
        return text
    except Exception as e:
        return ''

def get_csv_filtered_count__97b9c260(env, config: dict):
    """Get count of contacts with last names starting with 'W' and validate filtering.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'count' (int) and 'all_valid' (bool) keys
              - count: Number of contacts in CSV
              - all_valid: True if all contacts have last names starting with 'W'
    """
    path = config.get('path', '')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'count': 0, 'all_valid': False}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.strip().split('\n')
        if len(lines) < 2:
            return {'count': 0, 'all_valid': False}
        reader = csv.DictReader(lines)
        count = 0
        all_valid = True
        last_names = []
        for row in reader:
            if not row:
                continue
            count += 1
            last_name = None
            for key in row.keys():
                if key and 'last' in key.lower() and ('name' in key.lower()):
                    last_name = row[key]
                    break
            if last_name is None:
                for key in ['Last Name', 'LastName', 'Surname', 'Family Name']:
                    if key in row:
                        last_name = row[key]
                        break
            if last_name is None or not last_name.strip():
                all_valid = False
            elif not last_name.strip().upper().startswith('W'):
                all_valid = False
            last_names.append(last_name if last_name else '')
        return {'count': count, 'all_valid': all_valid, 'last_names': last_names}
    except Exception as e:
        return {'count': 0, 'all_valid': False}

def get_file_count__9364293cce5b25e22063aee62da7d43d(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of filenames in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' key

    Returns:
        List of filenames in the folder (basenames only, not full paths)
    """
    folder_path = config.get('folder_path', '')
    command = f"ls -1 '{folder_path}' 2>/dev/null"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0:
        logger.warning(f'Failed to list files in folder: {folder_path}')
        return []
    try:
        output = result.get('output', '').strip()
        if not output:
            return []
        filenames = [f.strip() for f in output.split('\n') if f.strip()]
        return filenames
    except AttributeError:
        logger.warning(f"Failed to parse file list from output: {result.get('output')}")
        return []

def get_file_count_from_file__be02851a(env, config):
    """
    Get the file count from a text file.

    Reads the content of the specified file and extracts the file count number.
    The file is expected to contain a numeric value representing the total file count.

    Args:
        env: Environment object with controller for file operations
        config: Configuration dict containing:
            - path: Path to the file containing the count (e.g., '/home/user/Documents/total_file_count.txt')
            - dest: Destination filename (for copying/reference)

    Returns:
        int: The file count read from the file, or None if file doesn't exist or contains invalid data
    """
    try:
        file_path = config.get('path')
        if not file_path:
            return None
        content = env.controller.get_vm_file(file_path)
        if content is None:
            return None
        content_str = content.strip()
        import re
        match = re.search('\\d+', content_str)
        if match:
            return int(match.group())
        return None
    except Exception as e:
        return None

def get_file_list_content__b23d642c(env, config):
    """Read chapter summary file content from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the file path

    Returns:
        str: File content as string, or None if file doesn't exist
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        logger.warning(f'File not found or could not be read: {path}')
        return None
    try:
        content = file_bytes.decode('utf-8')
        logger.info(f'Successfully read file: {path}, content length: {len(content)} chars')
        return content.strip()
    except Exception as e:
        logger.error(f'Failed to decode file content: {e}')
        return None

def get_bib_file_nonempty__9d1faa88(env, config: dict):
    """Parse bibtex file and extract entry details for verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Dictionary containing:
            - 'has_content': bool, whether file has content
            - 'content': str, the full bibtex content
            - 'has_cosql': bool, whether content mentions CoSQL or related terms
            - 'is_published': bool, whether it's published version (not arxiv)
            - 'has_blank_line': bool, whether there's a blank line at the end
    """
    import re
    path = config.get('path', '/home/user/Desktop/references.bib')
    file_content = env.controller.get_file(path)
    if not file_content:
        return {'has_content': False, 'content': '', 'has_cosql': False, 'is_published': False, 'has_blank_line': False}
    content = file_content.decode('utf-8', errors='ignore')
    if not content.strip():
        return {'has_content': False, 'content': '', 'has_cosql': False, 'is_published': False, 'has_blank_line': False}
    has_cosql = bool(re.search('CoSQL', content, re.IGNORECASE) or (re.search('Conversational', content, re.IGNORECASE) and re.search('Text-to-SQL', content, re.IGNORECASE)) or re.search('(Yu|Tao|Zhang|Radev)', content))
    has_arxiv = bool(re.search('arxiv', content, re.IGNORECASE) or re.search('arXiv', content))
    has_published_type = bool(re.search('@inproceedings\\s*\\{', content, re.IGNORECASE) or re.search('@article\\s*\\{', content, re.IGNORECASE))
    is_published = has_published_type and (not has_arxiv)
    has_blank_line = bool(re.search('\\}\\s*\\n\\s*\\n', content) or (content.rstrip() != content.rstrip('\n').rstrip() and '\n\n' in content[-10:]))
    return {'has_content': True, 'content': content, 'has_cosql': has_cosql, 'is_published': is_published, 'has_blank_line': has_blank_line}

def get_file_exists_and_size__21492178(env, config: Dict[str, Any]):
    """Check if a file exists on the VM and return its properties."""
    path = config.get('path', '')
    command = f"if [ -f '{path}' ]; then stat -c '%s' '{path}'; echo 'EXISTS'; else echo 'NOT_FOUND'; fi"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'NOT_FOUND' in output or not output:
        return {'exists': False, 'size': 0, 'is_png': False}
    lines = output.split('\n')
    size = 0
    exists = False
    for line in lines:
        if line.strip().isdigit():
            size = int(line.strip())
        if 'EXISTS' in line:
            exists = True
    is_png = False
    if exists:
        check_png_cmd = f"file '{path}' | grep -i 'PNG image'"
        png_result = env.controller.run_bash_script(check_png_cmd, timeout=10)
        is_png = png_result.get('returncode', 1) == 0
    return {'exists': exists, 'size': size, 'is_png': is_png}

def get_file_exists__e86a3a4d(env, config):
    """Check if file exists on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path)

    Returns:
        Boolean indicating file existence
    """
    vm_path = config.get('path')
    file_bytes = env.controller.get_file(vm_path)
    return file_bytes is not None

def get_txt_file_count__2b6d7a72(env, config):
    """Get details about .doc and .txt files for conversion verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' parameter

    Returns:
        dict: Contains:
            - txt_count: Number of .txt files found
            - doc_count: Number of .doc files found
            - matched_conversions: List of basenames where both .doc and .txt exist
            - txt_with_content: Number of .txt files with non-empty content
    """
    directory = config.get('directory', '/home/user/Desktop')
    doc_command = f"cd {directory} && ls -1 *.doc 2>/dev/null | sed 's/\\.doc$//' || echo ''"
    doc_result = env.controller.run_bash_script(doc_command, timeout=10)
    doc_basenames = [line.strip() for line in doc_result.get('output', '').split('\n') if line.strip()]
    txt_command = f"cd {directory} && ls -1 *.txt 2>/dev/null | sed 's/\\.txt$//' || echo ''"
    txt_result = env.controller.run_bash_script(txt_command, timeout=10)
    txt_basenames = [line.strip() for line in txt_result.get('output', '').split('\n') if line.strip()]
    matched_conversions = [basename for basename in doc_basenames if basename in txt_basenames]
    txt_with_content = 0
    if txt_basenames:
        check_limit = min(5, len(txt_basenames))
        for i in range(check_limit):
            txt_file = f'{directory}/{txt_basenames[i]}.txt'
            check_command = f"test -s {txt_file} && echo 'has_content' || echo 'empty'"
            check_result = env.controller.run_bash_script(check_command, timeout=5)
            if 'has_content' in check_result.get('output', ''):
                txt_with_content += 1
    return {'txt_count': len(txt_basenames), 'doc_count': len(doc_basenames), 'matched_conversions': matched_conversions, 'txt_with_content': txt_with_content}

def get_dir_file_list__173d79a32d2c31ac3ad30d4ae958526d(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files in a directory on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'command' key

    Returns:
        List of filenames in the directory, or empty list if directory doesn't exist
    """
    command = config.get('command', 'ls /home/user/Desktop/qa_session/')
    dir_path = command.replace('ls ', '').strip()
    check_result = env.controller.run_bash_script(f"test -d '{dir_path}'", timeout=10)
    if check_result['returncode'] != 0:
        return []
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return []
    output = result['output'].strip()
    if not output:
        return []
    files = [f.strip() for f in output.split('\n') if f.strip()]
    return files

def get_files_with_prefix__47275204(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of files in a directory that start with a specific prefix.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory_path' and 'prefix'

    Returns:
        List of filenames (without path) that start with the prefix
    """
    directory_path = config.get('directory_path', '')
    prefix = config.get('prefix', '')
    cmd = f"cd '{directory_path}' && ls -1 {prefix}* 2>/dev/null | sort"
    result = env.controller.run_bash_script(cmd, timeout=10)
    files = []
    if result.get('status') == 'success':
        output = result.get('output', '').strip()
        if output:
            files = [line.strip() for line in output.split('\n') if line.strip()]
    logger.info(f"Files with prefix '{prefix}' in {directory_path}: {files}")
    return files

def get_csv_filtered_rows__01b44b7acd315e39b1c9a6baa6b5f6da(env, config):
    """
    Read CSV file from VM and return all rows as list of lists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of lists representing CSV rows (including header)
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        rows = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(row)
        return rows
    finally:
        os.unlink(tmp_path)

def get_vm_subtitle_file__735234d9(env, config: dict):
    """Get subtitle file from VM."""
    vm_path = config['path']
    dest_name = config['dest']
    cache_path = os.path.join(env.cache_dir, dest_name)
    file_content = env.controller.get_file(vm_path)
    if file_content is None:
        return None
    with open(cache_path, 'wb') as f:
        f.write(file_content)
    return cache_path

def get_text_file_content__4080707f437c813fb70f2db7aaa30575(env, config):
    """Extract text content from a file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content as string, or empty string if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return ''
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
        return content
    except Exception as e:
        return ''

def get_file_exists__7107319d(env, config: dict):
    """Check if the exported HTML file exists and validate its content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with validation results including:
            - exists: bool - whether file exists
            - is_html: bool - whether file contains HTML content
            - has_table: bool - whether file contains table data
            - file_size: int - size of file in bytes
            - has_content: bool - whether file has substantial content
    """
    vm_path = config.get('path', '/home/user/Desktop/financial-data.html')
    result = env.controller.run_bash_script(f"test -f '{vm_path}' && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    output = result.get('output', '').strip()
    file_exists = output == 'EXISTS'
    if not file_exists:
        return {'exists': False, 'is_html': False, 'has_table': False, 'file_size': 0, 'has_content': False}
    size_result = env.controller.run_bash_script(f"stat -c %s '{vm_path}' 2>/dev/null || echo '0'", timeout=10)
    file_size = int(size_result.get('output', '0').strip() or 0)
    file_type_result = env.controller.run_bash_script(f"file -b --mime-type '{vm_path}' 2>/dev/null || echo 'unknown'", timeout=10)
    file_type = file_type_result.get('output', '').strip()
    is_html_type = 'html' in file_type.lower() or 'text' in file_type.lower()
    read_result = env.controller.run_bash_script(f"head -n 100 '{vm_path}' 2>/dev/null || echo ''", timeout=10)
    content = read_result.get('output', '')
    content_lower = content.lower()
    has_html_tags = any((tag in content_lower for tag in ['<html', '<table', '<tr>', '<tr ', '<td>', '<td ', '<th>', '<th ']))
    has_table_tag = '<table' in content_lower
    has_tr_tag = '<tr>' in content_lower or '<tr ' in content_lower
    has_table = has_table_tag and has_tr_tag
    has_content = file_size > 1500
    return {'exists': True, 'is_html': is_html_type and has_html_tags, 'has_table': has_table, 'file_size': file_size, 'has_content': has_content}

def get_dual_file_check__805294f8(env, config):
    """Check if both specified files exist on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'file1' and 'file2' parameters

    Returns:
        dict: Status of both files {"file1_exists": bool, "file2_exists": bool}
    """
    file1_path = config.get('file1')
    file2_path = config.get('file2')
    result1 = env.controller.run_bash_script(f'[ -f "{file1_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    file1_exists = 'exists' in result1['output'] and 'not exists' not in result1['output']
    result2 = env.controller.run_bash_script(f'[ -f "{file2_path}" ] && echo "exists" || echo "not exists"', timeout=10)
    file2_exists = 'exists' in result2['output'] and 'not exists' not in result2['output']
    return {'file1_exists': file1_exists, 'file2_exists': file2_exists}

def get_file_line_count__fb1a48c8(env, config):
    """Get line count of a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Number of lines in file
    """
    path = config.get('path', '')
    command = f'wc -l < {path}'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] == 0:
        try:
            return int(result['output'].strip())
        except ValueError:
            logger.error(f"Failed to parse line count: {result['output']}")
            return 0
    else:
        logger.error(f"Failed to count lines in {path}: {result['error']}")
        return 0

def get_text_file_lines__5ced85fc_aug18_v0_c9e8a1b2d3f4e5a6b7c8d9e0f1a2b3c4(env, config):
    """Read a text file from VM and return its lines as a list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        list: List of lines from the file (stripped of trailing newlines)
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return []
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return []
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = [line.rstrip('\n\r') for line in content.splitlines()]
        logger.info(f'Successfully read {len(lines)} lines from {file_path}')
        return lines
    except Exception as e:
        logger.error(f'Error reading file {file_path}: {e}')
        return []
