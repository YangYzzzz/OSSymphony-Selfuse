"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_python_class_definitions__8e31593051ddbcebea635607c768bb8c', 'get_python_score_display__ac3536a4', 'get_python_comments__f2c8b3d7', 'get_python_function_signatures__6082ba78', 'get_python_space_logic__f10704cc', 'get_python_functions__198be354', 'get_git_repo_exists__68f4a1f5', 'get_python_functions__9883d0c9e9f101970bd2e29a8a29d345', 'get_python_imports__261836e0618ec18a7a70e8da4837dfbf', 'get_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385', 'get_python_imports__5e897a929023bc43113113bb0ca8fb36', 'get_terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012', 'get_git_config__07fbe1c4', 'get_python_score_logic__84673dd6', 'get_python_pkg_check__a3b47e9754f6a01a17ed98f7d00d938c', 'get_python_imports__0a56ab11', 'get_git_status__763d7485', 'get_terminal_scrollback__8e9fd232f4307bd96cd64b8e7a1f9389', 'get_python_rotate_check__57357002', 'get_python_pause_feature__d5e5a07f', 'get_python_fall_speed_config__ca8e8c20', 'get_terminal_font_size__4d83e269ffe9a3de03f2bd05d7dacfe0', 'get_git_branch__400622aa', 'get_python_text_patterns__198be354', 'get_python_code_stats__198be354', 'get_python_hyperparams__526d49269c794c02e875c3ff44cdccae', 'get_python_docstrings__7a9e3b4c', 'get_python_imports__198be354', 'get_python_assignments__3e7d9c5f', 'get_python_lines_counter__20d8676e', 'get_python_gameover_screen__deb31e2f', 'get_python_grid_border__a01068b1', 'get_python_definitions__e8ec0313751ff65df15abd7d031a7c53', 'get_python_classes__198be354', 'get_bash_history__14fa116b50d2b7c945d2ede19b12f134', 'get_python_syntax__198be354', 'get_terminal_color_scheme__ad0133b72746f22f133946cc65aa7f55', 'get_git_repo_status__55495fb6b59196cc7cffae2b12e117ed', 'get_user_shell__49da98ad0a0985144d25f764fe852408', 'get_python_block_colors__b48e0411', 'get_git_repo_structure__f1a99656b1aa540fcb46d2aeba395a7e']

def get_python_class_definitions__8e31593051ddbcebea635607c768bb8c(env, config: dict):
    """Extract class definitions from Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Mapping of class names to their line ranges
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'classes': {}}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.split('\n')
        classes = {}
        current_class = None
        class_start = 0
        for (i, line) in enumerate(lines, 1):
            stripped = line.strip()
            if stripped.startswith('class '):
                if current_class:
                    classes[current_class] = {'start': class_start, 'end': i - 1}
                class_name = stripped.split('(')[0].replace('class ', '').replace(':', '').strip()
                current_class = class_name
                class_start = i
            elif current_class and len(line) > 0 and (line[0] not in [' ', '\t', '#']):
                if not stripped.startswith('class '):
                    classes[current_class] = {'start': class_start, 'end': i - 1}
                    current_class = None
        if current_class:
            classes[current_class] = {'start': class_start, 'end': len(lines)}
        return {'exists': True, 'classes': classes, 'class_count': len(classes)}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_python_score_display__ac3536a4(env, config):
    """
    Extract main.py to verify score display rendering is added.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of main.py
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/main.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    return content

def get_python_comments__f2c8b3d7(env, config: dict):
    """
    Extract comment lines from a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of comment lines, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = content.splitlines()
        comment_lines = [line.strip() for line in lines if line.strip().startswith('#')]
        return comment_lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_python_function_signatures__6082ba78(env, config: dict):
    """
    Extract function signatures (def statements) from a Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of function signature lines, or None if file doesn't exist
    """
    import re
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = content.splitlines()
        signatures = []
        for line in lines:
            stripped = line.strip()
            if re.match('^def\\s+\\w+', stripped):
                signatures.append(stripped)
        return signatures
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_python_space_logic__f10704cc(env, config):
    """
    Extract the go_space method from tetris.py to verify instant drop logic.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of the go_space method
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/tetris.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    lines = content.split('\n')
    in_method = False
    method_lines = []
    indent_level = None
    for line in lines:
        if 'def go_space(self)' in line:
            in_method = True
            indent_level = len(line) - len(line.lstrip())
            method_lines.append(line)
        elif in_method:
            if line.strip() and (not line.startswith(' ' * (indent_level + 1))):
                break
            method_lines.append(line)
    return '\n'.join(method_lines)

def get_python_functions__198be354(env, config):
    """Extract comprehensive Python code information from a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Information about functions, imports, classes, and code completeness
    """
    import re
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'function_count': 0, 'function_names': [], 'has_imports': False, 'import_count': 0, 'has_classes': False, 'class_count': 0, 'has_global_code': False}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    lines = content_str.splitlines()
    function_names = []
    imports = []
    classes = []
    has_global_code = False
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith('#'):
            continue
        if stripped.startswith('def '):
            match = re.match('def\\s+(\\w+)', stripped)
            if match:
                function_names.append(match.group(1))
        elif stripped.startswith('import ') or stripped.startswith('from '):
            imports.append(stripped)
        elif stripped.startswith('class '):
            match = re.match('class\\s+(\\w+)', stripped)
            if match:
                classes.append(match.group(1))
        elif '=' in stripped or any((stripped.startswith(kw) for kw in ['print(', 'if ', 'for ', 'while '])):
            has_global_code = True
    return {'function_count': len(function_names), 'function_names': function_names, 'has_imports': len(imports) > 0, 'import_count': len(imports), 'has_classes': len(classes) > 0, 'class_count': len(classes), 'has_global_code': has_global_code}

def get_git_repo_exists__68f4a1f5(env, config):
    """Check if a git repository exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' parameter

    Returns:
        dict: {'exists': bool, 'is_git_repo': bool}
    """
    repo_path = config.get('repo_path', '/home/user/transformers')
    command = f"test -d {repo_path} && test -d {repo_path}/.git && echo 'SUCCESS' || echo 'FAILED'"
    result = env.controller.run_bash_script(command, timeout=10)
    exists = 'SUCCESS' in result.get('output', '')
    return {'exists': exists, 'is_git_repo': exists}

def get_python_functions__9883d0c9e9f101970bd2e29a8a29d345(env, config: dict):
    """Extract function definitions from Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Function names and their signatures
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'functions': []}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.split('\n')
        functions = []
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('def '):
                func_sig = stripped.replace('def ', '').split(':')[0].strip()
                func_name = func_sig.split('(')[0].strip()
                is_method = len(line) > 0 and line[0] in [' ', '\t']
                functions.append({'name': func_name, 'signature': func_sig, 'is_method': is_method})
        methods = [f for f in functions if f['is_method']]
        top_level_funcs = [f for f in functions if not f['is_method']]
        return {'exists': True, 'all_functions': [f['name'] for f in functions], 'methods': [m['name'] for m in methods], 'top_level_functions': [f['name'] for f in top_level_funcs], 'total_count': len(functions), 'method_count': len(methods), 'function_count': len(top_level_funcs)}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_python_imports__261836e0618ec18a7a70e8da4837dfbf(env, config: dict):
    """Extract import statements from Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Categorized imports
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'exists': False, 'imports': []}
    try:
        content = file_bytes.decode('utf-8')
        lines = content.split('\n')
        standard_imports = []
        third_party_imports = []
        from_imports = []
        third_party_libs = ['torch', 'numpy', 'pandas', 'tensorflow', 'sklearn', 'matplotlib', 'scipy', 'requests', 'flask', 'django']
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('import '):
                module = stripped.replace('import ', '').split()[0].split('.')[0]
                if module in third_party_libs:
                    third_party_imports.append(stripped)
                else:
                    standard_imports.append(stripped)
            elif stripped.startswith('from '):
                parts = stripped.split()
                if len(parts) >= 2:
                    module_name = parts[1].split('.')[0]
                    if module_name in third_party_libs:
                        third_party_imports.append(stripped)
                        from_imports.append(stripped)
                    else:
                        standard_imports.append(stripped)
                        from_imports.append(stripped)
        return {'exists': True, 'standard_imports': standard_imports, 'third_party_imports': third_party_imports, 'from_imports': from_imports, 'total_imports': len(standard_imports) + len(third_party_imports)}
    except Exception as e:
        return {'exists': False, 'error': str(e)}

def get_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385(env, config):
    """Extract all Python comments from a file on VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the Python file on VM

    Returns:
        str: All comment lines, one per line
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        return ''
    comment_lines = []
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith('#'):
            comment_lines.append(line.rstrip())
        elif '#' in line:
            idx = line.index('#')
            comment_part = line[idx:].rstrip()
            if comment_part:
                comment_lines.append(comment_part)
    return '\n'.join(comment_lines)

def get_python_imports__5e897a929023bc43113113bb0ca8fb36(env, config):
    """Extract all unique import statements from a Python file on VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the Python file on VM

    Returns:
        str: All unique import statements, one per line
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        return ''
    import_lines = []
    seen_imports = set()
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith('import ') or stripped.startswith('from '):
            clean_line = stripped.split('#')[0].strip()
            if clean_line and clean_line not in seen_imports:
                seen_imports.add(clean_line)
                import_lines.append(clean_line)
    return '\n'.join(import_lines)

def get_terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012(env, config):
    """
    Get terminal cursor shape configuration from GNOME Terminal profile.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Terminal output containing cursor shape configuration
    """
    return env.controller.get_terminal_output()

def get_git_config__07fbe1c4(env, config):
    """Get git remote URL from repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' parameter

    Returns:
        str: Git remote URL or empty string if not found
    """
    repo_path = config.get('repo_path')
    command = f"cd {repo_path} && git config --get remote.origin.url 2>/dev/null || echo 'NOT_FOUND'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    return output if output != 'NOT_FOUND' else ''

def get_python_score_logic__84673dd6(env, config):
    """
    Extract the break_lines method from tetris.py to verify scoring logic.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of the break_lines method
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/tetris.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    lines = content.split('\n')
    in_method = False
    method_lines = []
    indent_level = None
    for line in lines:
        if 'def break_lines(self)' in line:
            in_method = True
            indent_level = len(line) - len(line.lstrip())
            method_lines.append(line)
        elif in_method:
            if line.strip() and (not line.startswith(' ' * (indent_level + 1))):
                break
            method_lines.append(line)
    return '\n'.join(method_lines)

def get_python_pkg_check__a3b47e9754f6a01a17ed98f7d00d938c(env, config: dict) -> str:
    """
    Execute pip list command to check for Python packages

    Args:
        env: DesktopEnv instance
        config: Configuration dict with command parameters

    Returns:
        str: Command output containing package list
    """
    command = config.get('command', [])
    result = env.controller.run_bash_script(' '.join(command) if isinstance(command, list) else command, timeout=30)
    output = result.get('output', '')
    logger.info(f'Python package check output: {output}')
    return output

def get_python_imports__0a56ab11(env, config: dict):
    """
    Extract import statements from a Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of import lines (strings), or None if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore')
        lines = content.splitlines()
        import_lines = []
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('import ') or stripped.startswith('from '):
                import_lines.append(stripped)
        return import_lines
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_git_status__763d7485(env, config: dict):
    """Check if a directory is a git repository.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Status info about git repository
    """
    path = config.get('path', '/home/user/Documents/Projects/OSWorld')
    result = env.controller.run_bash_script(f"cd '{path}' && git rev-parse --is-inside-work-tree 2>/dev/null", timeout=10)
    is_git_repo = result['returncode'] == 0 and result.get('output', '').strip() == 'true'
    branch_name = ''
    if is_git_repo:
        branch_result = env.controller.run_bash_script(f"cd '{path}' && git branch --show-current 2>/dev/null", timeout=10)
        if branch_result['returncode'] == 0:
            branch_name = branch_result.get('output', '').strip()
    return {'is_git_repo': is_git_repo, 'branch': branch_name}

def get_terminal_scrollback__8e9fd232f4307bd96cd64b8e7a1f9389(env, config):
    """
    Get terminal scrollback lines configuration from GNOME Terminal profile.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Terminal output containing scrollback configuration
    """
    return env.controller.get_terminal_output()

def get_python_rotate_check__57357002(env, config):
    """
    Extract the rotate method from tetris.py to verify collision checking is added.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of the rotate method
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/tetris.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    lines = content.split('\n')
    in_rotate = False
    rotate_lines = []
    indent_level = None
    for line in lines:
        if 'def rotate(self)' in line:
            in_rotate = True
            indent_level = len(line) - len(line.lstrip())
            rotate_lines.append(line)
        elif in_rotate:
            if line.strip() and (not line.startswith(' ' * (indent_level + 1))):
                break
            rotate_lines.append(line)
    return '\n'.join(rotate_lines)

def get_python_pause_feature__d5e5a07f(env, config):
    """
    Extract code from main.py and tetris.py to verify pause feature is added.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        dict: Contains main.py and tetris.py relevant sections
    """
    main_path = config.get('main_path', '/home/user/Desktop/tetris/main.py')
    tetris_path = config.get('tetris_path', '/home/user/Desktop/tetris/tetris.py')
    result = {'main_py': '', 'tetris_py': ''}
    main_bytes = env.controller.get_file(main_path)
    if main_bytes:
        result['main_py'] = main_bytes.decode('utf-8')
    tetris_bytes = env.controller.get_file(tetris_path)
    if tetris_bytes:
        result['tetris_py'] = tetris_bytes.decode('utf-8')
    return result

def get_python_fall_speed_config__ca8e8c20(env, config):
    """
    Extract settings.py and main.py to verify fall speed is configurable.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        dict: Contains settings.py and main.py content
    """
    settings_path = config.get('settings_path', '/home/user/Desktop/tetris/settings.py')
    main_path = config.get('main_path', '/home/user/Desktop/tetris/main.py')
    result = {'settings_py': '', 'main_py': ''}
    settings_bytes = env.controller.get_file(settings_path)
    if settings_bytes:
        result['settings_py'] = settings_bytes.decode('utf-8')
    main_bytes = env.controller.get_file(main_path)
    if main_bytes:
        result['main_py'] = main_bytes.decode('utf-8')
    return result

def get_terminal_font_size__4d83e269ffe9a3de03f2bd05d7dacfe0(env, config):
    """
    Get terminal font size configuration from GNOME Terminal profile.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Terminal output containing font size
    """
    return env.controller.get_terminal_output()

def get_git_branch__400622aa(env, config):
    """Get current git branch name.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' parameter

    Returns:
        str: Current branch name or empty string if not a git repo
    """
    repo_path = config.get('repo_path')
    command = f"cd {repo_path} && git branch --show-current 2>/dev/null || echo 'NOT_GIT_REPO'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    return output if output != 'NOT_GIT_REPO' else ''

def get_python_text_patterns__198be354(env, config):
    """Search for specific text patterns in a Python file and validate code structure.

    This function extracts the Python file and checks for:
    - Presence of required text patterns (key code components)
    - Valid Python syntax
    - Proper code structure (imports, class definitions)
    - Substantial content length

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'patterns' keys

    Returns:
        dict: Pattern matching results with code validity check
    """
    file_path = config.get('path', '')
    patterns_to_find = config.get('patterns', [])
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'pattern_matches': {}, 'is_valid_python': False, 'has_imports': False, 'has_class_def': False, 'content_length': 0}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    pattern_matches = {}
    for pattern in patterns_to_find:
        pattern_matches[pattern] = pattern in content_str
    is_valid_python = False
    try:
        ast.parse(content_str)
        is_valid_python = True
    except:
        is_valid_python = False
    has_imports = 'import ' in content_str or 'from ' in content_str
    has_class_def = 'class BigramLanguageModel' in content_str
    content_length = len(content_str.strip())
    return {'exists': True, 'pattern_matches': pattern_matches, 'is_valid_python': is_valid_python, 'has_imports': has_imports, 'has_class_def': has_class_def, 'content_length': content_length}

def get_python_code_stats__198be354(env, config):
    """Get statistics about Python code vs comments.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Code statistics (total_lines, code_lines, comment_lines, blank_lines)
    """
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'total_lines': 0, 'code_lines': 0, 'comment_lines': 0, 'blank_lines': 0}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    lines = content_str.splitlines()
    code_lines = 0
    comment_lines = 0
    blank_lines = 0
    for line in lines:
        stripped = line.strip()
        if not stripped:
            blank_lines += 1
        elif stripped.startswith('#'):
            comment_lines += 1
        else:
            code_lines += 1
    return {'exists': True, 'total_lines': len(lines), 'code_lines': code_lines, 'comment_lines': comment_lines, 'blank_lines': blank_lines}

def get_python_hyperparams__526d49269c794c02e875c3ff44cdccae(env, config):
    """Extract hyperparameter assignments from a Python file on VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the Python file on VM

    Returns:
        str: Hyperparameter assignments, one per line
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        return ''
    hyperparam_keywords = ['batch_size', 'block_size', 'learning_rate', 'lr', 'max_iters', 'eval_interval', 'device', 'n_embd', 'n_head', 'n_layer', 'dropout', 'eval_iters', 'epochs', 'hidden_size', 'num_layers', 'vocab_size', 'max_length', 'temperature', 'momentum']
    hyperparam_lines = []
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped or stripped.startswith('#'):
            continue
        for keyword in hyperparam_keywords:
            if stripped.startswith(f'{keyword} =') or stripped.startswith(f'{keyword}='):
                hyperparam_lines.append(line.rstrip())
                break
    return '\n'.join(hyperparam_lines)

def get_python_docstrings__7a9e3b4c(env, config: dict):
    """
    Extract docstrings from a text file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: File content, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore')
        return content
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_python_imports__198be354(env, config):
    """Extract import statements and code metrics from a Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Information about imports, line count, and function definitions
    """
    import re
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'import_count': 0, 'unique_modules': [], 'total_lines': 0, 'non_empty_lines': 0, 'function_count': 0, 'class_count': 0}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    lines = content_str.splitlines()
    import_lines = []
    unique_modules = set()
    function_count = 0
    class_count = 0
    non_empty_lines = 0
    for line in lines:
        stripped = line.strip()
        if stripped and (not stripped.startswith('#')):
            non_empty_lines += 1
        if stripped.startswith('import ') or stripped.startswith('from '):
            import_lines.append(stripped)
            if stripped.startswith('import '):
                parts = stripped[7:].split()
                if parts:
                    module = parts[0].split(',')[0].strip()
                    unique_modules.add(module)
            elif stripped.startswith('from '):
                parts = stripped[5:].split()
                if parts and len(parts) > 1:
                    module = parts[0].strip()
                    unique_modules.add(module)
        if stripped.startswith('def '):
            function_count += 1
        elif stripped.startswith('class '):
            class_count += 1
    return {'import_count': len(import_lines), 'unique_modules': sorted(list(unique_modules)), 'total_lines': len(lines), 'non_empty_lines': non_empty_lines, 'function_count': function_count, 'class_count': class_count}

def get_python_assignments__3e7d9c5f(env, config: dict):
    """
    Extract numeric assignment statements from a Python file using AST parsing.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of assignment lines with numeric values, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    if not file_path:
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return None
        content = file_bytes.decode('utf-8', errors='ignore')
        try:
            tree = ast.parse(content)
            assignments = []
            lines = content.splitlines()
            for node in ast.walk(tree):
                if isinstance(node, ast.Assign):
                    if _is_numeric_value(node.value):
                        assignment_str = _extract_assignment_source(node, lines, content)
                        if assignment_str:
                            assignments.append(assignment_str)
                elif isinstance(node, ast.AugAssign):
                    if _is_numeric_value(node.value):
                        assignment_str = _extract_assignment_source(node, lines, content)
                        if assignment_str:
                            assignments.append(assignment_str)
                elif isinstance(node, ast.AnnAssign) and node.value:
                    if _is_numeric_value(node.value):
                        assignment_str = _extract_assignment_source(node, lines, content)
                        if assignment_str:
                            assignments.append(assignment_str)
            return assignments if assignments else None
        except SyntaxError:
            return _fallback_extraction(content)
    except Exception as e:
        print(f'Error reading file: {e}')
        return None

def get_python_lines_counter__20d8676e(env, config):
    """
    Extract tetris.py and main.py to verify lines cleared counter tracking and display.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        dict: Contains tetris.py and main.py content
    """
    tetris_path = config.get('tetris_path', '/home/user/Desktop/tetris/tetris.py')
    main_path = config.get('main_path', '/home/user/Desktop/tetris/main.py')
    result = {'tetris_py': '', 'main_py': ''}
    tetris_bytes = env.controller.get_file(tetris_path)
    if tetris_bytes:
        result['tetris_py'] = tetris_bytes.decode('utf-8')
    main_bytes = env.controller.get_file(main_path)
    if main_bytes:
        result['main_py'] = main_bytes.decode('utf-8')
    return result

def get_python_gameover_screen__deb31e2f(env, config):
    """
    Extract main.py to verify game over screen and restart logic are added.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of main.py
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/main.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    return content

def get_python_grid_border__a01068b1(env, config):
    """
    Extract main.py to verify grid border rendering is enhanced.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of main.py
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/main.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    return content

def get_python_definitions__e8ec0313751ff65df15abd7d031a7c53(env, config):
    """Extract all function and class definitions from a Python file on VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the Python file on VM

    Returns:
        str: All function and class definitions, one per line
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return ''
    try:
        content = file_bytes.decode('utf-8')
    except UnicodeDecodeError:
        return ''
    definition_lines = []
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith('def ') or stripped.startswith('class '):
            definition_lines.append(line.rstrip())
    return '\n'.join(definition_lines)

def get_python_classes__198be354(env, config):
    """Extract class definitions from a Python file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Information about classes (count, class_names)
    """
    import re
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'class_count': 0, 'class_names': []}
    try:
        content_str = file_content.decode('utf-8')
    except:
        content_str = file_content.decode('utf-8', errors='ignore')
    lines = content_str.splitlines()
    class_names = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith('class '):
            match = re.match('class\\s+(\\w+)', stripped)
            if match:
                class_names.append(match.group(1))
    return {'class_count': len(class_names), 'class_names': class_names}

def get_bash_history__14fa116b50d2b7c945d2ede19b12f134(env, config):
    """Get bash history content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: Content of bash history
    """
    command = 'cat ~/.bash_history'
    result = env.controller.run_bash_script(command, timeout=10)
    if result['returncode'] == 0:
        return result['output']
    return ''

def get_python_syntax__198be354(env, config):
    """Extract code from Colab notebook and check if saved file matches and has valid syntax.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Syntax validation and extraction verification results
    """
    import ast
    import logging
    from playwright.sync_api import sync_playwright
    logger = logging.getLogger('desktopenv.getters.python_syntax')
    file_path = config.get('path', '')
    file_content = env.controller.get_file(file_path)
    if not file_content:
        return {'exists': False, 'valid_syntax': False, 'extracted_from_notebook': False, 'saved_code': '', 'notebook_code': '', 'error': 'File not found'}
    try:
        saved_code = file_content.decode('utf-8')
    except Exception as e:
        return {'exists': True, 'valid_syntax': False, 'extracted_from_notebook': False, 'saved_code': '', 'notebook_code': '', 'error': f'Decode error: {str(e)}'}
    valid_syntax = False
    try:
        ast.parse(saved_code)
        valid_syntax = True
    except SyntaxError as e:
        pass
    except Exception as e:
        pass
    notebook_code = ''
    extraction_error = None
    try:
        host = env.vm_ip
        port = env.chromium_port
        remote_debugging_url = f'http://{host}:{port}'
        with sync_playwright() as p:
            try:
                browser = p.chromium.connect_over_cdp(remote_debugging_url)
                logger.info('[COLAB_EXTRACT] Connected to Chrome instance')
            except Exception as e:
                logger.warning(f'[COLAB_EXTRACT] Failed to connect to Chrome: {e}')
                extraction_error = f'Failed to connect to Chrome: {e}'
                browser = None
            if browser:
                colab_page = None
                for context in browser.contexts:
                    for page in context.pages:
                        try:
                            page_url = page.url
                            if 'colab.research.google.com' in page_url:
                                colab_page = page
                                logger.info(f'[COLAB_EXTRACT] Found Colab page: {page_url}')
                                break
                        except Exception as e:
                            continue
                    if colab_page:
                        break
                if colab_page:
                    try:
                        colab_page.wait_for_load_state('domcontentloaded', timeout=10000)
                        code_cells = []
                        selectors_to_try = ['.CodeMirror-code', '.inputarea', '[data-test-id="code-cell-editor"]', '.code-cell', 'colab-cell[type="code"]', '.code']
                        for selector in selectors_to_try:
                            try:
                                elements = colab_page.query_selector_all(selector)
                                if elements:
                                    logger.info(f"[COLAB_EXTRACT] Found {len(elements)} elements with selector '{selector}'")
                                    for elem in elements:
                                        text = elem.text_content()
                                        if text and text.strip():
                                            code_cells.append(text.strip())
                                    if code_cells:
                                        break
                            except Exception as e:
                                logger.debug(f"[COLAB_EXTRACT] Selector '{selector}' failed: {e}")
                                continue
                        if not code_cells:
                            logger.info('[COLAB_EXTRACT] Trying to extract code from page content')
                            try:
                                pre_elements = colab_page.query_selector_all('pre, code')
                                for elem in pre_elements:
                                    text = elem.text_content()
                                    if text and text.strip() and (len(text.strip()) > 10):
                                        code_cells.append(text.strip())
                            except Exception as e:
                                logger.debug(f'[COLAB_EXTRACT] Failed to extract from pre/code tags: {e}')
                        if code_cells:
                            notebook_code = '\n'.join(code_cells)
                            logger.info(f'[COLAB_EXTRACT] Extracted {len(code_cells)} code cells, total {len(notebook_code)} characters')
                        else:
                            logger.warning('[COLAB_EXTRACT] No code cells found in notebook')
                            extraction_error = 'No code cells found in notebook'
                    except Exception as e:
                        logger.error(f'[COLAB_EXTRACT] Error extracting code from Colab: {e}')
                        extraction_error = f'Error extracting code: {e}'
                else:
                    logger.warning('[COLAB_EXTRACT] Could not find Colab page')
                    extraction_error = 'Colab page not found'
                browser.close()
    except Exception as e:
        logger.error(f'[COLAB_EXTRACT] Error during Playwright operation: {e}')
        extraction_error = f'Playwright error: {e}'
    extracted_from_notebook = False
    if notebook_code and saved_code:
        saved_normalized = ''.join(saved_code.split())
        notebook_normalized = ''.join(notebook_code.split())
        if len(saved_normalized) > 0:
            overlap_ratio = 0
            if len(notebook_normalized) > 0:
                common_length = sum((1 for (c1, c2) in zip(saved_normalized, notebook_normalized) if c1 == c2))
                overlap_ratio = common_length / max(len(saved_normalized), len(notebook_normalized))
                if saved_normalized in notebook_normalized:
                    overlap_ratio = 1.0
            extracted_from_notebook = overlap_ratio > 0.3
            logger.info(f'[COLAB_EXTRACT] Overlap ratio: {overlap_ratio:.2f}, extracted_from_notebook: {extracted_from_notebook}')
    return {'exists': True, 'valid_syntax': valid_syntax, 'extracted_from_notebook': extracted_from_notebook, 'saved_code': saved_code, 'notebook_code': notebook_code, 'error': extraction_error}

def get_terminal_color_scheme__ad0133b72746f22f133946cc65aa7f55(env, config):
    """
    Get terminal color scheme configuration from GNOME Terminal profile.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: Terminal output containing color scheme info
    """
    return env.controller.get_terminal_output()

def get_git_repo_status__55495fb6b59196cc7cffae2b12e117ed(env, config: dict):
    """Check if a git repository exists and get its remote URL.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        dict: Status information with keys:
            - 'exists': bool, whether directory exists
            - 'is_git_repo': bool, whether it's a valid git repository
            - 'remote_url': str, git remote URL if available
    """
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    repo_path = config['repo_path']
    result = {'exists': False, 'is_git_repo': False, 'remote_url': ''}
    check_dir_cmd = f"test -d {repo_path} && echo 'exists' || echo 'not_exists'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': check_dir_cmd, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        if output == 'exists':
            result['exists'] = True
        else:
            return result
    else:
        logger.error('Failed to check directory existence. Status code: %d', response.status_code)
        return result
    check_git_cmd = f"test -d {repo_path}/.git && echo 'is_git' || echo 'not_git'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': check_git_cmd, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        if output == 'is_git':
            result['is_git_repo'] = True
        else:
            return result
    else:
        logger.error('Failed to check git repository. Status code: %d', response.status_code)
        return result
    get_url_cmd = f'cd {repo_path} && git remote get-url origin'
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': get_url_cmd, 'shell': True})
    if response.status_code == 200:
        result['remote_url'] = response.json()['output'].strip()
    else:
        logger.error('Failed to get git remote URL. Status code: %d', response.status_code)
    return result

def get_user_shell__49da98ad0a0985144d25f764fe852408(env, config):
    """Get user information including shell, home directory, and password status.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'username' key

    Returns:
        dict: Dictionary with 'username', 'shell', 'home', and 'password_set' keys,
            or None if user doesn't exist
    """
    username = config.get('username', '')
    import requests
    vm_ip = env.vm_ip
    port = env.server_port
    command = f"getent passwd {username} | cut -d':' -f1,6,7 || echo 'USER_NOT_FOUND'"
    response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
    if response.status_code == 200:
        output = response.json()['output'].strip()
        if output == 'USER_NOT_FOUND' or not output:
            return None
        parts = output.split(':')
        if len(parts) < 3:
            logger.error('Unexpected passwd output format: %s', output)
            return None
        user_name = parts[0]
        home_dir = parts[1]
        shell = parts[2]
        shadow_command = f"sudo getent shadow {username} | grep -q '^{username}:' && echo 'PASSWORD_EXISTS' || echo 'NO_PASSWORD'"
        shadow_response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': shadow_command, 'shell': True})
        password_set = False
        if shadow_response.status_code == 200:
            shadow_output = shadow_response.json()['output'].strip()
            password_set = shadow_output == 'PASSWORD_EXISTS'
        return {'username': user_name, 'home': home_dir, 'shell': shell, 'password_set': password_set}
    else:
        logger.error('Failed to get user info. Status code: %d', response.status_code)
        return None

def get_python_block_colors__b48e0411(env, config):
    """
    Extract the shapes definition section from block.py to verify colors are assigned.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with file paths

    Returns:
        str: Content of block.py focusing on Block class and shapes
    """
    file_path = config.get('path', '/home/user/Desktop/tetris/block.py')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return ''
    content = file_bytes.decode('utf-8')
    return content

def get_git_repo_structure__f1a99656b1aa540fcb46d2aeba395a7e(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a git repository exists and get its structure.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'repo_path' key

    Returns:
        Dict with 'exists', 'is_git_repo', 'remote_url', and 'files' keys
    """
    repo_path = config.get('repo_path', '')
    result = {'exists': False, 'is_git_repo': False, 'remote_url': '', 'files': []}
    check_dir_cmd = f"[ -d '{repo_path}' ] && echo 'EXISTS' || echo 'NOT_EXISTS'"
    dir_check = env.controller.run_bash_script(check_dir_cmd, timeout=10)
    if dir_check.get('output', '').strip() == 'EXISTS':
        result['exists'] = True
        git_check_cmd = f"[ -d '{repo_path}/.git' ] && echo 'GIT_REPO' || echo 'NOT_GIT'"
        git_check = env.controller.run_bash_script(git_check_cmd, timeout=10)
        if git_check.get('output', '').strip() == 'GIT_REPO':
            result['is_git_repo'] = True
            remote_url_cmd = f"git -C '{repo_path}' remote get-url origin 2>/dev/null || echo ''"
            remote_result = env.controller.run_bash_script(remote_url_cmd, timeout=10)
            if remote_result.get('returncode') == 0:
                remote_url = remote_result.get('output', '').strip()
                result['remote_url'] = remote_url
        ls_cmd = f"cd '{repo_path}' && find . -maxdepth 2 -type f -o -type d | sort"
        ls_result = env.controller.run_bash_script(ls_cmd, timeout=30)
        if ls_result.get('returncode') == 0:
            files = ls_result.get('output', '').strip().split('\n')
            result['files'] = [f.strip() for f in files if f.strip()]
    logger.info(f'Git repo check for {repo_path}: {result}')
    return result
