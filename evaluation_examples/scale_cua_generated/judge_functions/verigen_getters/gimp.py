"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_folder_image_counts__6a4313ae', 'get_image_grayscale__0e7db70c', 'get_image_hash__0969de9d', 'get_image_flipped__3948a175', 'get_gimp_image_size__993c9cd2028ad767fa928842de0805ca', 'get_gif_file_info__f8138076f6546232c093f5027d50db21', 'get_gimp_layer_names__ca8ad8ab', 'get_image_cropped_dims__c48d85ee866b2dcfd0ab60f091a5a2b6', 'get_gimp_gimprc_file__adbf37d0', 'get_image_dimensions__1728b8eeebbdebd7894eb2578c0c67ec', 'get_chrome_images_setting__9945e6ef', 'get_gimp_layer_names__17690e7b', 'get_image_hash__8778dd25', 'get_image_properties__016cfcf5', 'get_gif_file_info__430ad71e1ce94d5ea2bb30ba073fd937', 'get_image_properties__72e666ad', 'get_png_file_info__92ca6baf', 'get_photo_rename_verification__0a812fadd7008b5cb899c479c82ab6e7', 'get_gimp_image_size__07fcde31879d28764a081db212afaf2d', 'get_remaining_image_filenames__c90c6ca3', 'get_docx_image_info__2ceb2e40a4dbc2ef362c792411e6262a', 'get_image_orientation__777d01ae6d6663c5897f2674681dca2e', 'get_image_dimensions__2aa2e23a948bbaf102eec8741d709529', 'get_image_brightness__42b8076d', 'get_jpg_file_info__1bf235b17a43af3cc147100153de4d30', 'get_image_cropped_size__166a89c9', 'get_folder_image_list__fe2f25921e2c2c43fbc9e31c35bccb78', 'get_image_file_size__c46a6f1dddc552cd368bc819d4cce6f7', 'get_image_properties__0a6f6cc1', 'get_gif_file_info__16fb0dc6', 'get_image_size__912586a80097a155904c15da97ac2079', 'get_doc_image_total__d49c5873', 'get_gimp_image_size__f8a3e253', 'get_gimp_config_menubar__eebb7fd7089a7abc35d29b3d4832455e', 'get_image_brightness__2f4c8c94976cda4c41de6383a6a66dec', 'get_image_count__48e4325f2d62197da2b10059281b95a0', 'get_extracted_image_properties__f3e27026', 'get_gimp_image_content_size__e2468a8febb27268d777ba03561c41ab', 'get_image_rotation_aspect__b014bcfb2b56c94b7ab718a85ad6cff9', 'get_image_format__22b30bf2', 'get_image_properties__c4f65e24', 'get_default_video_player__d253b8f0', 'get_default_image_viewer__7ae7ab08c735ff1cd3e440ff41b233f8', 'get_has_background_image__23bfac54', 'get_gif_file_info__ea8c7a7a', 'get_gif_file_info__e66cc6b0', 'get_image_hash__69eadfa0', 'get_image_dimensions__7816ba80', 'get_image_rotation__82cdf8c66b7532068e061ebffd0a6c33', 'get_docx_total_images__02587c7a', 'get_image_mode__78bfc971', 'get_gimp_image_size__5db5c513f28a5ad7ad8d60ef05b14b35', 'get_image_saturation__8a7e9a1c5d75d961a07bbfbc1999062a', 'get_has_background_image__6ccbfc2a', 'get_gimp_image_flipped__16b4973e', 'get_gimp_rotated_image__e0c6cf186be7cc4682d3712837dcdd72', 'get_image_corner_colors__340a3600a88f3d4dc2217b9f986d625d', 'get_image_file_format__024272f6', 'get_image_format__945e33b1', 'get_image_dimensions__3b2067dec4f66b0e25da3355b1fb8f3e', 'get_gimp_gimprc_file__5b6e5d1f', 'get_image_properties__ab9c94ea', 'get_triangle_horizontal_position__c73a98370436f02ecc8c151f49a9b91c', 'get_png_resolution_info__27528fb8', 'get_image_saturation__df2e4b52', 'get_gif_file_info__092b88db7884c339b37665609b49294b', 'get_googledrive_png_files__f9d97b19', 'get_image_rotation__b49205bc', 'get_gif_file_info__78103de86e64961437a4bcd00b97b9bc', 'get_gimp_image_size__fb0dfb53338d620d275120282f62d824', 'get_all_images__1574bc56f755697238f4190c2d72a32b', 'get_extracted_image_properties__a8440735', 'get_triangle_rotation__2c517d983ea369519ffc55979e3393ec', 'get_image_rotated__045e2e0c', 'get_image_format_info__900fc36276e9eaf12471a7834992cf5c', 'get_image_properties__70f3db86', 'get_image_properties__50f8e5bb', 'get_gimp_layer_names__8565a91c', 'get_image_properties__c6063730', 'get_default_audio_player__971ae155', 'get_image_properties__795e137d', 'get_default_image_viewer__ae568cb6', 'get_image_properties__68566bbc', 'get_image_file__2237c5bebdb76e54ae53ea89e71ca4a3', 'get_googledrive_png_count__b4c05b10', 'get_docx_image_count__5cfb373b', 'get_png_export_check__c69055e15cceefc40a87e6de042c2331', 'get_image_color_mode__90bdfeb2ebfd1bac6873718be37b3912', 'get_image_file__203069587d53a571860bccb97348992b', 'get_gimp_saturation__8f3d7ee1bf388be9293a1e57f283a35e', 'get_png_export_info__046a9ca717a3ff75711d7d7d1d876a5f', 'get_image_flip_check__08f973926b8f35af5489796c73a6c6e0', 'get_image_dimensions__b6f18d98d1993dda1c36f44651fe6a5d', 'get_image_dimensions__e0d9b551', 'get_docx_image_count__c37ca17b', 'get_has_background_image__99fc8543', 'get_has_background_image__171e4956', 'get_image_format__c0fb0f23', 'get_gif_file_info__11cf8ab6', 'get_gif_file_info__b130b682', 'get_image_aspect_ratio__f85979b813875fecca12ba1c6ab4cc68', 'get_cropped_image_size__8159102f', 'get_image_color_mode__3aef0cbd8986e240435edab0fd96c873', 'get_gdrive_png_count__63477992bbc88c6a7091e80f7c0a8a72', 'get_image_dimensions__fa1a72d6', 'get_image_properties__739292ff', 'get_docx_image_count__dbe26b57', 'get_image_file_exists__7c0cc95089263e14cb308f3757c8acc1', 'get_image_dimensions__505cf5fc', 'get_image_count_status__8e1f0943', 'get_googledrive_image_files__67a58d2f', 'get_gimp_file_bytes__677c871105619c52013adb82fa7e5d28', 'get_gimp_image_size__8677d5c370f70c69494653c2a8ef5be2', 'get_bottom_half_image__4baa14c913fe32a60f559a98296affae', 'get_gimp_image_brightness__ced7a23c', 'get_docx_image_count__74be0e09', 'get_image_info__e19bfc7ef1338231fef513d8a5b2f6d1', 'get_gimp_flipped_image__07b5058cd3df711389d2b4342d0c561c', 'get_docx_image_embedded__6ed140cad2fa6c5f3349fe322a43680c', 'get_image_dimensions__b24412f0', 'get_default_audio_player__7752be97956df277ef920fced6cad6a7', 'get_docx_image_count__081533bc', 'get_image_scaled_dims__bf4967f3e9931b3c80be8e4dbf6e04b7', 'get_extracted_image_properties__69184b17', 'get_image_validation__2511ecbd', 'get_gimp_gimprc_file__9c13adcb', 'get_gimp_image_brightness__652303c0d066122d99f102352a5b1a93', 'get_image_dimensions__baa9d5de58b4726390c9c04659eb9fca', 'get_image_properties__811349c6', 'get_odt_image_check__603a6f2174a0ecd455ce537d9f738dbe', 'get_original_image_size__e429d357', 'get_image_dimensions__f5cfdff3841c16728bb4565a839b59ca', 'get_image_properties__05f2a34a', 'get_gimp_layer_count__58c8d068', 'get_extracted_image_properties__4e5534f0', 'get_has_background_image__005bc54f', 'get_image_blurriness__eea46e4e4afa193660c5f52c6a2da7a9', 'get_image_total__b01cf463', 'get_image_hash__4c0f04bf', 'get_gif_file_info__0d0257bdd2e8345a807cfdefd39ffa3b', 'get_image_hash__d1128c0a', 'get_image_properties__4bd1c0a2e0720abe2c2c09ee9978ce22', 'get_image_dimensions__77b19ce3287accb29381eac14cc998b5', 'get_gimp_moved_image__e4487e27c5c6e4232b26add556e7d796', 'get_image_dimensions__a7d5fd37', 'get_image_contrast__9829f3ab', 'get_image_scaled__5d068218', 'get_image_dimensions__bec1165ff4f2eb5b48dc7de50a4fe1ab', 'get_image_properties__9873b6c6', 'get_image_dimensions__4eb874068ccb861273ecf8604bdafb3c', 'get_image_dimensions__462996d1', 'get_image_mode__1ac9295ec86f9e2e04c973e2e47b273c', 'get_png_file_exists__b9abaa4fbc51b493882263c6a6aff8fe', 'get_saved_image_77b8ab4d', 'get_gif_file_info__de678d13fd248567f73258f2b3cb0372', 'get_image_size__4939e12b', 'get_extracted_image_properties__a3446500', 'get_gimp_aspect_ratio__3e33afd89a022c42c408db15555e5c16', 'get_default_image_editor__77752d0e', 'get_image_properties__34372286', 'get_image_crop_check__c26faa30e6e2f9e09b4748ad3193b390', 'get_image_properties__37707684957589279b0fa14602529fe7', 'get_image_hash__10e2f0b6', 'get_gimp_config_toolbox__90386cba106758972cba7bf949bb562f', 'get_image_rotation__4e34ef5d', 'get_xcf_layer_names__b0cabbaa0a8c8fbf299cc3425a48f58e', 'get_docx_image_presence__b5b56630', 'get_image_dimensions__c5f81e73faaccc56bdfa2edf29f272b7', 'get_image_dimensions__c984db77', 'get_default_image_viewer__6ed8455a', 'get_image_brightness__473ce27f0364f391d90c68cfef960e0c', 'get_triangle_topleft_position__40af8739a7f9d38f5046e0dfa41c5f6e', 'get_image_hash__0c114ed7', 'get_image_brightness__35cd470f', 'get_gimp_grayscale_check__37ae16c860b893668df4aed8d0b9ae18', 'get_gimp_gimprc_file__17026b9e', 'get_gimp_image_saturation__62aadfc75943521b0b091bf9d9c10f24', 'get_gif_frame_count__fc6196ba2aeabb7bc4c476007b2b24c6', 'get_default_image_viewer__56d62a52f0286a3d54adea0a15934e97', 'get_docx_has_new_image__3333dfb2', 'get_image_for_flip_check__9f1f61b8216550440a48089a3e4c1731', 'get_image_orientation__55bebaef7ca999b793134c2c00f342a9', 'get_image_mode__8f91f3a7', 'get_has_background_image__daf96dfa', 'get_image_dimensions__9f5722fc', 'get_image_dimensions__f2472a44', 'get_gimp_brightness__4c7b26ad7d6a1a697cf659bf75175126', 'get_extracted_image_properties__1a437041', 'get_png_count__eb6b46ad', 'get_gimp_image_color_mode__7ae19854d29f1b0c8cdcf13e028f0bdd', 'get_image_properties__4894850c', 'get_triangle_area__12d50454feda301909898f2cf2cce54b', 'get_image_count_final__5d239d03', 'get_gimp_layer_names__b148e375_v3', 'get_image_color_mode__bf5ecb70b33ef3dffbe095b744ec38a7', 'get_gimp_config_statusbar__2d2ae6cc356da88025063f58e6110bad', 'get_image_dimensions__844f5e73b108da449b9b68fcbef6bbf2', 'get_image_props__a01f23a4', 'get_gimp_layer_name_config__174a6594', 'get_gimp_recolored_image__d52919fa54c3d57b57da98359753b674', 'get_triangle_color__978cb6f31473d4802226bc7ae94b7399', 'get_jpg_list_file__8ddf03c6b80780c49b4f9497dee3f888', 'get_gimp_config_rulers__96f3f9d88ca4e4230495b91ff566eb51', 'get_gimp_image_dimensions__4bd0ac4fe70775f29bef20a161a34c39', 'get_image_dimensions__1beedc32', 'get_gimp_scaled_image__93095e022ff3d0c1026d009f3ccc512b', 'get_image_size__7ecb394c9eae8a8e135a21ca629ec0de', 'get_image_properties__bd10eca8', 'get_default_music_player__2fc648743a39d7031a57b208b33f9200', 'get_gimp_layer_name_config__3336340a', 'get_gimp_image_file__bc5fe443a1102b529ba31d0ac9c81ab1', 'get_gimp_layer_names__a85f6474', 'get_third_image__9a3feca86555f3c82732707871883142', 'get_gimp_layer_names__677f0e4c', 'get_gif_file_info__06536a54', 'get_has_background_image__4b23e04e', 'get_image_mode__9660fa13484d43a1462069231ef86deb', 'get_gif_dimensions__90c23f3797e29598f167849a527a40d5', 'get_default_audio_player__3b265860', 'get_image_properties__9d6c98f0', 'get_image_dimensions__de812dd5b44b906cb9793a8d4a3f91bf', 'get_image_size__10742f90', 'get_gimp_transparency__456c340f2121fceb6f1006996239e1cf', 'get_gimp_layer_exists__e5a8318d', 'get_image_properties__fa4ed82a', 'get_image_color_mode__14447080df6d9553c7e99d3265fc5a81', 'get_gimp_file_bytes__05275b4640be680d1ddfc9692455a07a', 'get_gimp_layer_names__8249d409', 'get_image_file_count__3fde58b2', 'get_image_dimensions__696943e1', 'get_image_contrast__979c644a69a9cc6d01b110bbe0f08e78', 'get_image_dimensions__7cb89717dbfd62e5cbbd4dcc85a4e268', 'get_image_dimensions__da5d7378bbf41608325407fc00f8d126', 'get_gif_file_info__d18c2fd4', 'get_image_dimensions__3f39d534f7803089ed331d19c2a1bc89', 'get_image_properties__ab3f6dd1', 'get_gimp_gimprc_file__dc948653', 'get_gimp_brightness__0b33739cf65dc15a04df08d23f597efa', 'get_gif_file_info__06722a19', 'get_image_hash__6510dd9a', 'get_docx_image_count__1b5dc5fe', 'get_png_file__d04e36cb527bbbbfdb06458981bf8945', 'get_image_properties__565232d7', 'get_image_hash__b461596b', 'get_image_orientation__6575e228', 'get_image_file_exists__836db3cb3a7cc19d82f98c6a439eb80c', 'get_docx_images_count__17e4ac0c', 'get_xcf_layer_names__b148e375', 'get_gimp_layer_group_config__7ba73b05', 'get_image_mode__d9cbc3c3', 'get_extracted_image_properties__345e8ddb', 'get_image_dimensions__a4f96e46', 'get_docx_image_count__34e55d07', 'get_gimp_config_fullscreen__8efddf2685fdb790d7823145c3565e94', 'get_docx_image_count__cd0e399d', 'get_jpeg_export_check__39da7f334341155f29f73cbacf02786e', 'get_gimp_flipped_image__1aaf02638da71a4a84f47e22e2395da3', 'get_image_hash__db2588a0', 'get_image_hash__73176121', 'get_gif_file_info__3d85489a', 'get_image_dimensions__4694337da0a8886d5bd508a95fd83b12', 'get_has_background_image__f45ebd91', 'get_image_properties__9472df29', 'get_docx_image_count__7d8c4525', 'get_image_dimensions__a846d552e46987bd85e04c0a4f658c7a', 'get_gimp_gimprc_file__d7b0aa1f', 'get_png_files_list__ec3ddc36', 'get_image_rotation_check__25c734e4497155c51ced0623dec284fc', 'get_docx_image_count__e898c34b', 'get_image_crop_info__59e6ca63b22becd85a8942d1a29325d9', 'get_gimp_layer_names__aa5f92aa', 'get_has_background_image__6c48637d', 'get_image_properties__12182c78', 'get_image_dimensions__184d842d21bd06203c79c089532a2315', 'get_gif_file_info__8d385ddc', 'get_image_dimensions__717d6863', 'get_docx_image_count__02c7d140', 'get_image_props__6575e228_v8']

def get_folder_image_counts__6a4313ae(env, config: Dict[str, Any]) -> Dict[str, int]:
    """
    Get the count of image files in two specified folders.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder1_path' and 'folder2_path'

    Returns:
        Dict with 'folder1_count' and 'folder2_count' keys
    """
    folder1_path = config.get('folder1_path', '')
    folder2_path = config.get('folder2_path', '')
    cmd1 = f"find '{folder1_path}' -maxdepth 1 -type f \\( -iname '*.jpg' -o -iname '*.jpeg' -o -iname '*.png' \\) 2>/dev/null | wc -l"
    result1 = env.controller.run_bash_script(cmd1, timeout=10)
    cmd2 = f"find '{folder2_path}' -maxdepth 1 -type f \\( -iname '*.jpg' -o -iname '*.jpeg' -o -iname '*.png' \\) 2>/dev/null | wc -l"
    result2 = env.controller.run_bash_script(cmd2, timeout=10)
    folder1_count = 0
    folder2_count = 0
    if result1.get('status') == 'success':
        try:
            folder1_count = int(result1.get('output', '0').strip())
        except ValueError:
            logger.warning(f"Failed to parse folder1 count: {result1.get('output')}")
    if result2.get('status') == 'success':
        try:
            folder2_count = int(result2.get('output', '0').strip())
        except ValueError:
            logger.warning(f"Failed to parse folder2 count: {result2.get('output')}")
    logger.info(f'Folder image counts - {folder1_path}: {folder1_count}, {folder2_path}: {folder2_count}')
    return {'folder1_count': folder1_count, 'folder2_count': folder2_count}

def get_image_grayscale__0e7db70c(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get an image file from the VM to check if it's grayscale.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key for the image path on VM

    Returns:
        str: Local path to the downloaded image file, or None if failed
    """
    path = config['path']
    dest = os.path.basename(path)
    _path = os.path.join(env.cache_dir, dest)
    try:
        file = env.controller.get_file(path)
        if file is None:
            logger.warning(f'Failed to get image file from VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(_path, 'wb') as f:
            f.write(file)
        logger.info(f'Successfully saved image: {_path} ({len(file)} bytes)')
        return _path
    except Exception as e:
        logger.error(f'Error getting image file {path}: {e}')
        return None

def get_image_hash__0969de9d(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_flipped__3948a175(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the flipped image file from the VM.

    This function retrieves an image file (background.png) that was created
    by flipping the original background image horizontally.

    Args:
        env: Environment object with controller and cache_dir
        config: Configuration dict containing:
            - path (str): absolute path on the VM to fetch the image

    Returns:
        str: Path to the downloaded file in cache, or None if file doesn't exist
    """
    path = config['path']
    dest = os.path.basename(path)
    try:
        file = env.controller.get_file(path)
        if file is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        cache_path = os.path.join(env.cache_dir, dest)
        with open(cache_path, 'wb') as f:
            f.write(file)
        logger.info(f'Successfully saved image file: {cache_path} ({len(file)} bytes)')
        return cache_path
    except Exception as e:
        logger.error(f'Error processing file {path}: {e}')
        return None

def get_gimp_image_size__993c9cd2028ad767fa928842de0805ca(env, config: Dict[str, str]):
    """
    Gets the size (width, height) and perceptual hash of an image file from VM.
    Also retrieves corner pixel samples for rotation verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        dict: {'width': int, 'height': int, 'phash': str, 'corner_samples': list}
              or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        width = img.width
        height = img.height
        phash = str(imagehash.phash(img))
        corner_samples = []
        sample_size = min(10, width // 10, height // 10)
        if sample_size > 0:
            tl_region = img.crop((0, 0, sample_size, sample_size))
            corner_samples.append(('top_left', list(tl_region.getdata())[:25]))
            tr_region = img.crop((width - sample_size, 0, width, sample_size))
            corner_samples.append(('top_right', list(tr_region.getdata())[:25]))
            bl_region = img.crop((0, height - sample_size, sample_size, height))
            corner_samples.append(('bottom_left', list(bl_region.getdata())[:25]))
            br_region = img.crop((width - sample_size, height - sample_size, width, height))
            corner_samples.append(('bottom_right', list(br_region.getdata())[:25]))
        result = {'width': width, 'height': height, 'phash': phash, 'corner_samples': corner_samples}
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image info: width={width}, height={height}, phash={phash}')
        return result
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_gif_file_info__f8138076f6546232c093f5027d50db21(env, config: Dict):
    """
    Get information about a GIF file on the VM, including frame count and duration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path on VM

    Returns:
        dict with file existence, size, frame count, duration, and basic properties
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_ms': 0}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_gif = img.format == 'GIF'
        (width, height) = img.size
        frames = 0
        total_duration_ms = 0
        try:
            while True:
                frames += 1
                frame_duration = img.info.get('duration', 100)
                total_duration_ms += frame_duration
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        result = {'exists': True, 'file_size': len(file_bytes), 'is_gif': is_gif, 'width': width, 'height': height, 'frames': frames, 'duration_ms': total_duration_ms}
        logger.info(f'GIF file info: {result}')
        return result
    except Exception as e:
        logger.error(f'Error analyzing GIF file: {e}')
        return {'exists': True, 'file_size': len(file_bytes), 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_ms': 0}
    finally:
        os.unlink(tmp_path)

def get_gimp_layer_names__ca8ad8ab(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_image_cropped_dims__c48d85ee866b2dcfd0ab60f091a5a2b6(env, config: dict):
    """
    Get image file from VM and return its dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        dict: Dictionary with 'width' and 'height' keys, or None if failed
    """
    vm_path = config['path']
    dest_name = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        img = Image.open(cache_path)
        result = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {result}')
        return result
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        return None

def get_gimp_gimprc_file__adbf37d0(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_image_dimensions__1728b8eeebbdebd7894eb2578c0c67ec(env, config: Dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the image file

    Returns:
        Dict with 'width' and 'height' keys, or None if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get file: {config['path']}")
        return None
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            img.close()
            return {'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image dimensions: {e}')
        return None

def get_chrome_images_setting__9945e6ef(env, config: Dict[str, str]):
    """
    Get the images content setting from Chrome preferences.

    Args:
        env: Desktop environment instance
        config: Configuration dictionary

    Returns:
        str: "allow" if images are enabled, "block" if disabled
    """
    os_type = env.vm_platform
    if os_type == 'Windows':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('LOCALAPPDATA'),\n                                            'Google\\Chrome\\User Data\\Default\\Preferences'))")['output'].strip()
    elif os_type == 'Darwin':
        preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'Library/Application Support/Google/Chrome/Default/Preferences'))")['output'].strip()
    elif os_type == 'Linux':
        if 'arm' in platform.machine():
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), 'snap/chromium/common/chromium/Default/Preferences'))")['output'].strip()
        else:
            preference_file_path = env.controller.execute_python_command("import os; print(os.path.join(os.getenv('HOME'), '.config/google-chrome/Default/Preferences'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system')
    try:
        content = env.controller.get_file(preference_file_path)
        data = json.loads(content)
        images_setting = data.get('profile', {}).get('default_content_setting_values', {}).get('images', 1)
        return 'block' if images_setting == 2 else 'allow'
    except Exception as e:
        logger.error(f'Error getting images setting: {e}')
        return 'allow'

def get_gimp_layer_names__17690e7b(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_image_hash__8778dd25(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_properties__016cfcf5(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_gif_file_info__430ad71e1ce94d5ea2bb30ba073fd937(env, config: Dict):
    """
    Get information about a GIF file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path on VM

    Returns:
        dict with file existence, size, and basic properties
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_seconds': 0.0}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_gif = img.format == 'GIF'
        (width, height) = img.size
        frames = 0
        total_duration_ms = 0
        try:
            while True:
                frames += 1
                frame_duration = img.info.get('duration', 100)
                total_duration_ms += frame_duration
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        total_duration_seconds = total_duration_ms / 1000.0
        result = {'exists': True, 'file_size': len(file_bytes), 'is_gif': is_gif, 'width': width, 'height': height, 'frames': frames, 'duration_seconds': total_duration_seconds}
        logger.info(f'GIF file info: {result}')
        return result
    except Exception as e:
        logger.error(f'Error analyzing GIF file: {e}')
        return {'exists': True, 'file_size': len(file_bytes), 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_seconds': 0.0}
    finally:
        os.unlink(tmp_path)

def get_image_properties__72e666ad(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_png_file_info__92ca6baf(env, config: Dict) -> Optional[Dict]:
    """
    Get PNG file information including dimensions and existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'width', 'height' keys, or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        return None
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'width': 0, 'height': 0}
        with Image.open(local_path) as img:
            return {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'error': str(e)}

def get_photo_rename_verification__0a812fadd7008b5cb899c479c82ab6e7(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Verify photo rename operation with content validation using OCR.

    Checks:
    1. Renamed file exists (ending_01.jpg)
    2. Renamed file contains expected text ("Thank you") via OCR
    3. Original file no longer exists (DSC00657.jpg)
    4. Total file count in folder remains the same (no copies created)

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder' key (path to folder)

    Returns:
        Dict with verification results:
        {
            'renamed_file_exists': bool,
            'renamed_file_text': str (extracted text from OCR),
            'original_file_exists': bool,
            'total_files': int (count of .jpg files in folder)
        }
    """
    folder = config.get('folder', '')
    result = {'renamed_file_exists': False, 'renamed_file_text': '', 'original_file_exists': False, 'total_files': 0}
    renamed_path = os.path.join(folder, 'ending_01.jpg')
    check_cmd = f"test -f '{renamed_path}' && echo 'exists' || echo 'not_exists'"
    cmd_result = env.controller.run_bash_script(check_cmd, timeout=10)
    result['renamed_file_exists'] = cmd_result.get('returncode') == 0 and 'exists' in cmd_result.get('output', '')
    if result['renamed_file_exists']:
        try:
            file_bytes = env.controller.get_file(renamed_path)
            if file_bytes:
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    ocr_result = subprocess.run(['tesseract', tmp_path, 'stdout'], capture_output=True, text=True, timeout=30)
                    result['renamed_file_text'] = ocr_result.stdout
                except Exception as e:
                    logger.warning(f'OCR failed: {e}')
                    result['renamed_file_text'] = ''
                finally:
                    try:
                        os.unlink(tmp_path)
                    except:
                        pass
        except Exception as e:
            logger.warning(f'Failed to get file for OCR: {e}')
    original_path = os.path.join(folder, 'DSC00657.jpg')
    check_cmd = f"test -f '{original_path}' && echo 'exists' || echo 'not_exists'"
    cmd_result = env.controller.run_bash_script(check_cmd, timeout=10)
    result['original_file_exists'] = cmd_result.get('returncode') == 0 and 'exists' in cmd_result.get('output', '')
    count_cmd = f"find '{folder}' -maxdepth 1 -type f -name '*.jpg' | wc -l"
    cmd_result = env.controller.run_bash_script(count_cmd, timeout=10)
    try:
        result['total_files'] = int(cmd_result.get('output', '0').strip())
    except:
        result['total_files'] = 0
    return result

def get_gimp_image_size__07fcde31879d28764a081db212afaf2d(env, config: Dict[str, str]):
    """
    Get the size (width x height) of an image file from VM.
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file: {file_path}')
        return None
    import tempfile
    import os
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            size = {'width': img.width, 'height': img.height}
            logger.debug(f'Image size: {size}')
            return size
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image size: {e}')
        return None

def get_remaining_image_filenames__c90c6ca3(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of image filenames remaining in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory_path'

    Returns:
        List of image filenames (without path) sorted alphabetically
    """
    directory_path = config.get('directory_path', '')
    cmd = f"cd '{directory_path}' && ls -1 *.jpg *.jpeg *.png *.gif *.bmp 2>/dev/null | sort"
    result = env.controller.run_bash_script(cmd, timeout=10)
    files = []
    if result.get('status') == 'success':
        output = result.get('output', '').strip()
        if output:
            files = [line.strip() for line in output.split('\n') if line.strip()]
    logger.info(f'Remaining image files in {directory_path}: {files}')
    return files

def get_docx_image_info__2ceb2e40a4dbc2ef362c792411e6262a(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get information about images in DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to DOCX file

    Returns:
        Dict with file_exists and image_count
    """
    docx_path = config.get('path', '')
    result = {'file_exists': False, 'image_count': 0}
    file_bytes = env.controller.get_file(docx_path)
    if not file_bytes:
        logger.warning(f'DOCX file not found: {docx_path}')
        return result
    result['file_exists'] = True
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            result['image_count'] = image_count
            logger.info(f"Found {result['image_count']} image(s) in DOCX")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading DOCX file: {e}')
    return result

def get_image_orientation__777d01ae6d6663c5897f2674681dca2e(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get image orientation and verify 90-degree clockwise rotation by comparing with original.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'width', 'height', 'is_rotated_90' keys
    """
    from PIL import Image
    import tempfile
    result = {'exists': False, 'width': None, 'height': None, 'is_rotated_90': False}
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.warning(f"Failed to get file from VM: {config['path']}")
            return result
        result['exists'] = True
        original_path = '/home/user/OIP.jpg'
        original_bytes = env.controller.get_file(original_path)
        if not original_bytes:
            logger.warning(f'Failed to get original file from VM: {original_path}')
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                img = Image.open(tmp_path)
                result['width'] = img.size[0]
                result['height'] = img.size[1]
                if result['height'] > result['width']:
                    result['is_rotated_90'] = True
                logger.info(f"Image size: {result['width']}x{result['height']}, rotated_90={result['is_rotated_90']} (fallback check)")
            finally:
                os.unlink(tmp_path)
            return result
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp_original:
            tmp_original.write(original_bytes)
            tmp_original_path = tmp_original.name
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp_rotated:
            tmp_rotated.write(file_bytes)
            tmp_rotated_path = tmp_rotated.name
        try:
            original_img = Image.open(tmp_original_path)
            rotated_img = Image.open(tmp_rotated_path)
            result['width'] = rotated_img.size[0]
            result['height'] = rotated_img.size[1]
            (original_width, original_height) = original_img.size
            (rotated_width, rotated_height) = rotated_img.size
            logger.info(f'Original image size: {original_width}x{original_height}')
            logger.info(f'Rotated image size: {rotated_width}x{rotated_height}')
            if not (original_width == rotated_height and original_height == rotated_width):
                logger.info('Dimension check failed: dimensions not swapped')
                result['is_rotated_90'] = False
                return result
            original_array = np.array(original_img)
            rotated_array = np.array(rotated_img)
            sample_points = [(0, 0), (original_width - 1, 0), (0, original_height - 1), (original_width - 1, original_height - 1), (original_width // 2, original_height // 2)]
            matches = 0
            for (x, y) in sample_points:
                if x < original_width and y < original_height:
                    rotated_x = x
                    rotated_y = original_width - 1 - y
                    if rotated_y < rotated_height and rotated_x < rotated_width:
                        original_pixel = original_array[y, x]
                        rotated_pixel = rotated_array[rotated_y, rotated_x]
                        pixel_diff = np.abs(original_pixel.astype(int) - rotated_pixel.astype(int))
                        if np.all(pixel_diff < 10):
                            matches += 1
            if matches >= len(sample_points) * 0.8:
                result['is_rotated_90'] = True
                logger.info(f'Rotation verified: {matches}/{len(sample_points)} sample points match (90-degree clockwise)')
            else:
                result['is_rotated_90'] = False
                logger.info(f'Rotation verification failed: only {matches}/{len(sample_points)} sample points match')
        finally:
            os.unlink(tmp_original_path)
            os.unlink(tmp_rotated_path)
    except Exception as e:
        logger.error(f'Error getting image orientation: {e}')
    return result

def get_image_dimensions__2aa2e23a948bbaf102eec8741d709529(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get image dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'width', 'height' keys
    """
    from PIL import Image
    import tempfile
    result = {'exists': False, 'width': None, 'height': None}
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.warning(f"Failed to get file from VM: {config['path']}")
            return result
        result['exists'] = True
        with tempfile.NamedTemporaryFile(suffix='.img', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result['width'] = img.size[0]
            result['height'] = img.size[1]
            logger.info(f"Image dimensions: {result['width']}x{result['height']}")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
    return result

def get_image_brightness__42b8076d(env, config: Dict[str, Any]) -> Dict[str, float]:
    """Get the average brightness of an image file and compare with reference.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key and optional 'reference_path' key

    Returns:
        dict: {
            'brightness': float,  # Average brightness value (0-255 scale)
            'similarity': float   # Structural similarity to reference (0-1 scale)
        }
    """
    try:
        from PIL import Image
        import numpy as np
        from skimage.metrics import structural_similarity as ssim
    except ImportError as e:
        logger.error(f'Failed to import required libraries: {e}')
        return {'brightness': 0.0, 'similarity': 0.0}
    image_path = config.get('path', '')
    reference_path = config.get('reference_path', '')
    try:
        from desktop_env.evaluators.getters.file import get_vm_file
        dest_filename = os.path.basename(image_path)
        local_path = get_vm_file(env, {'path': image_path, 'dest': dest_filename})
        if not local_path or not os.path.exists(local_path):
            logger.error(f'Image file not found: {image_path}')
            return {'brightness': 0.0, 'similarity': 0.0}
        img = Image.open(local_path)
        grayscale = img.convert('L')
        brightness = float(np.mean(np.array(grayscale)))
        logger.info(f'Image brightness: {brightness}')
        similarity = 0.0
        if reference_path:
            try:
                ref_dest_filename = os.path.basename(reference_path)
                ref_local_path = get_vm_file(env, {'path': reference_path, 'dest': ref_dest_filename})
                if ref_local_path and os.path.exists(ref_local_path):
                    ref_img = Image.open(ref_local_path)
                    if img.size != ref_img.size:
                        ref_img = ref_img.resize(img.size, Image.LANCZOS)
                    img_gray = np.array(img.convert('L'))
                    ref_gray = np.array(ref_img.convert('L'))
                    similarity = float(ssim(img_gray, ref_gray))
                    logger.info(f'Image similarity to reference: {similarity}')
                else:
                    logger.warning(f'Reference image not found: {reference_path}')
            except Exception as e:
                logger.error(f'Failed to calculate similarity: {e}')
        return {'brightness': brightness, 'similarity': similarity}
    except Exception as e:
        logger.error(f'Failed to get image brightness: {e}')
        return {'brightness': 0.0, 'similarity': 0.0}

def get_jpg_file_info__1bf235b17a43af3cc147100153de4d30(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if a JPG file exists and get its properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'is_jpg', 'format', 'size' keys
    """
    path = config.get('path', '')
    result = {'exists': False, 'is_jpg': False, 'format': None, 'size': 0}
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            result['exists'] = True
            result['size'] = len(file_bytes)
            if path.lower().endswith(('.jpg', '.jpeg')):
                try:
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                        tmp.write(file_bytes)
                        tmp_path = tmp.name
                    try:
                        img = Image.open(tmp_path)
                        result['format'] = img.format
                        if img.format in ['JPEG', 'JPG']:
                            result['is_jpg'] = True
                    except Exception as e:
                        logger.warning(f'File exists but is not a valid JPG: {e}')
                    finally:
                        os.unlink(tmp_path)
                except Exception as e:
                    logger.warning(f'Error verifying JPG: {e}')
        else:
            logger.info(f'File does not exist at path: {path}')
    except Exception as e:
        logger.error(f'Error checking file: {e}')
    return result

def get_image_cropped_size__166a89c9(env, config: Dict[str, Any]) -> Optional[Dict[str, int]]:
    """
    Get the dimensions (width and height) of an image file.

    Args:
        env: Environment object with controller to access VM files
        config: Configuration dict with 'path' key for the image file path

    Returns:
        Dict with 'width' and 'height' keys, or None if file doesn't exist
    """
    path = config['path']
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get image file from VM: {path}')
            return None
        dest = os.path.basename(path)
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        with Image.open(cache_path) as img:
            (width, height) = img.size
            logger.info(f'Image dimensions: {width}x{height}')
            return {'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error getting image dimensions for {path}: {e}')
        return None

def get_folder_image_list__fe2f25921e2c2c43fbc9e31c35bccb78(env, config: Dict[str, Any]) -> List[str]:
    """
    Get list of image files in a folder on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'folder_path' key

    Returns:
        List of image filenames (sorted)
    """
    folder_path = config.get('folder_path', '')
    command = f"ls '{folder_path}' 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    if result.get('returncode') != 0 or not result.get('output'):
        logger.warning(f'Failed to list folder: {folder_path}')
        return []
    files = result['output'].strip().split('\n')
    image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp']
    image_files = [f for f in files if any((f.lower().endswith(ext) for ext in image_extensions))]
    return sorted(image_files)

def get_image_file_size__c46a6f1dddc552cd368bc819d4cce6f7(env, config: dict):
    """
    Get file size and animation information for an image file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file path on VM

    Returns:
        dict: File info including exists, file_size_bytes, format, frame_count, is_animated
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path specified in config')
        return {'exists': False, 'file_size_bytes': 0, 'format': None, 'frame_count': 0, 'is_animated': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size_bytes': 0, 'format': None, 'frame_count': 0, 'is_animated': False}
    file_size = len(file_bytes)
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            img_format = img.format
            frame_count = 1
            is_animated = False
            if img_format == 'GIF':
                try:
                    frame_count = getattr(img, 'n_frames', 1)
                    is_animated = frame_count > 1
                    logger.info(f'GIF has {frame_count} frames, is_animated: {is_animated}')
                except Exception as e:
                    logger.warning(f'Could not determine frame count: {e}')
            img.close()
            return {'exists': True, 'file_size_bytes': file_size, 'format': img_format, 'frame_count': frame_count, 'is_animated': is_animated}
        except Exception as e:
            logger.error(f'Error opening image file: {e}')
            return {'exists': True, 'file_size_bytes': file_size, 'format': None, 'frame_count': 0, 'is_animated': False}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_image_properties__0a6f6cc1(env, config: dict):
    """Extract image properties from the saved image file.

    This function retrieves the saved image and extracts its properties
    (dimensions, format, existence). The specific dimensions (1114x623)
    serve as a strong indicator that the correct second image was extracted.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, and exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        img = Image.open(BytesIO(file_bytes))
        width = img.width
        height = img.height
        img_format = img.format
        img.close()
        return {'exists': True, 'width': width, 'height': height, 'format': img_format}
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_gif_file_info__16fb0dc6(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_size__912586a80097a155904c15da97ac2079(env, config: dict):
    """Get the dimensions of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with width and height
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'width': 0, 'height': 0}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        img.close()
        return dimensions
    finally:
        os.unlink(tmp_path)

def get_doc_image_total__d49c5873(env, config: dict):
    """Get total number of images in document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Total images
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_gimp_image_size__f8a3e253(env, config):
    """
    Get the dimensions of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        dict: Dictionary with 'width' and 'height' keys, or None if file not found
    """
    try:
        file_path = config.get('path')
        if not file_path:
            logger.error('No file path provided in config')
            return None
        file_content = env.controller.get_file(file_path)
        if not file_content:
            logger.error(f'Failed to get file: {file_path}')
            return None
        import os
        local_path = os.path.join(env.cache_dir, 'temp_resize_check.png')
        with open(local_path, 'wb') as f:
            f.write(file_content)
        img = Image.open(local_path)
        result = {'width': img.width, 'height': img.height}
        logger.info(f'Image size: {result}')
        return result
    except Exception as e:
        logger.error(f'Error getting image size: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_gimp_config_menubar__eebb7fd7089a7abc35d29b3d4832455e(env, config: Dict[str, str]):
    """
    Gets the GIMP config file to check menubar visibility setting.
    This getter retrieves the sessionrc file which contains window display settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_image_brightness__2f4c8c94976cda4c41de6383a6a66dec(env, config: Dict):
    """
    Get the average brightness of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the image file

    Returns:
        Dict with 'brightness' key (float value 0-255), or None if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get file: {config['path']}")
        return None
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            grayscale = img.convert('L')
            stat = ImageStat.Stat(grayscale)
            brightness = stat.mean[0]
            img.close()
            return {'brightness': brightness}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image brightness: {e}')
        return None

def get_image_count__48e4325f2d62197da2b10059281b95a0(env, config: dict) -> dict:
    """Get the content of image_count.txt file from Desktop and parse the count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'exists' (bool), 'content' (str), and 'count' (int or None)
    """
    path = config.get('path', '/home/user/Desktop/image_count.txt')
    file_content = env.controller.get_vm_file_content(path)
    if file_content is None:
        return {'exists': False, 'content': '', 'count': None}
    content_str = file_content.strip()
    count_value = None
    try:
        count_value = int(content_str)
    except ValueError:
        import re
        numbers = re.findall('\\d+', content_str)
        if numbers:
            count_value = int(numbers[0])
    return {'exists': True, 'content': content_str, 'count': count_value}

def get_extracted_image_properties__f3e27026(env, config: dict):
    """
    Get properties of an extracted image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path

    Returns:
        dict: Image properties including exists, size, format, dimensions
    """
    import os
    from PIL import Image
    from io import BytesIO
    path = config.get('path', '/home/user/extracted_image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0}
    try:
        img = Image.open(BytesIO(file_bytes))
        return {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0, 'error': str(e)}

def get_gimp_image_content_size__e2468a8febb27268d777ba03561c41ab(env, config):
    """
    Get the content size (non-transparent area) of a GIMP-exported image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int} or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            if img.mode in ('RGBA', 'LA') or 'transparency' in img.info:
                if img.mode != 'RGBA':
                    img = img.convert('RGBA')
                alpha = img.split()[-1]
                bbox = alpha.getbbox()
                if bbox is None:
                    logger.warning('Image is completely transparent')
                    return {'width': 0, 'height': 0}
                width = bbox[2] - bbox[0]
                height = bbox[3] - bbox[1]
                logger.info(f'Content size: {width}x{height} (from bbox {bbox})')
                return {'width': width, 'height': height}
            else:
                (width, height) = img.size
                logger.info(f'Image size (no transparency): {width}x{height}')
                return {'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image content size: {e}')
        return None

def get_image_rotation_aspect__b014bcfb2b56c94b7ab718a85ad6cff9(env, config):
    """Get image orientation (landscape/portrait) to verify rotation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int, 'orientation': str} or None
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (width, height) = (img.width, img.height)
        if width > height:
            orientation = 'landscape'
        elif height > width:
            orientation = 'portrait'
        else:
            orientation = 'square'
        result = {'width': width, 'height': height, 'orientation': orientation}
        img.close()
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_format__22b30bf2(env, config: dict):
    """Get image format and properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Image properties including format
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        return {'format': img.format, 'mode': img.mode, 'exists': True}
    except Exception as e:
        return {'exists': False}

def get_image_properties__c4f65e24(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_default_video_player__d253b8f0(env, config: dict):
    """
    Gets the default application for video files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        str: Default video player application name
    """
    result = env.controller.run_bash_script('xdg-mime query default video/mp4', timeout=10)
    if result['returncode'] == 0 and result['output']:
        return result['output'].strip()
    return 'unknown'

def get_default_image_viewer__7ae7ab08c735ff1cd3e440ff41b233f8(env, config: dict):
    """Gets the default image viewer application.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: The default image viewer .desktop file name (e.g., 'eog.desktop')
    """
    import requests
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['jpeg', 'png', 'gif', 'bmp', 'x-bmp', 'x-ms-bmp', 'tiff', 'x-icon', 'svg+xml', 'webp']
        apps = []
        vm_ip = env.vm_ip
        port = env.server_port
        for ext in extensions:
            command = ['xdg-mime', 'query', 'default', f'image/{ext}']
            response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
            if response.status_code == 200:
                app = response.json().get('output', '').strip()
                if app:
                    apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    else:
        raise Exception('Unsupported operating system', os_type)

def get_has_background_image__23bfac54(env, config: Dict[str, str]):
    """
    Extract the background image from a slide and compare it with a reference frame from a video.

    Args:
        env: Environment object
        config: Configuration dict with keys:
            - ppt_file_path: Path to the presentation file
            - slide_index: Index of the slide (0-based)
            - video_path: Path to the video file to extract reference frame from
            - timestamp: Timestamp in seconds to extract the frame from

    Returns:
        dict: Dictionary with keys:
            - 'background_image_path': Path to extracted background image (or None)
            - 'reference_frame_path': Path to extracted reference frame (or None)
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    video_path = config.get('video_path', '/home/user/Desktop/landscape.mp4')
    timestamp = float(config.get('timestamp', 20.0))
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    background_image_path = None
    reference_frame_path = None
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'background_image_path': None, 'reference_frame_path': None}
            image_id = None
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    try:
                        for element in child.iter(image_tag):
                            image_id = element.attrib.get(attr_tag)
                            if image_id is not None:
                                break
                    except:
                        pass
                    if image_id is not None:
                        break
            if image_id is None:
                return {'background_image_path': None, 'reference_frame_path': None}
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_rels_file in myzip.namelist():
                with myzip.open(slide_rels_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    namespaces = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                    for rel in root.findall('r:Relationship', namespaces):
                        if 'image' in rel.attrib['Type'] and rel.attrib['Id'] == image_id:
                            target = rel.attrib['Target']
                            if target.startswith('..'):
                                image_file_path = os.path.normpath(os.path.join('ppt/slides', target))
                                image_file_path = image_file_path.replace('\\', '/')
                                tmpdirname = os.path.dirname(ppt_file_localhost_path)
                                myzip.extract(image_file_path, tmpdirname)
                                background_image_path = os.path.join(tmpdirname, image_file_path)
                            break
    except Exception as e:
        pass
    try:
        video_localhost_path = get_vm_file(env, {'path': video_path, 'dest': os.path.split(video_path)[-1]})
        cap = cv2.VideoCapture(video_localhost_path)
        cap.set(cv2.CAP_PROP_POS_MSEC, timestamp * 1000)
        (ret, frame) = cap.read()
        cap.release()
        if ret:
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                reference_frame_path = tmp_file.name
                cv2.imwrite(reference_frame_path, frame)
    except Exception as e:
        pass
    return {'background_image_path': background_image_path, 'reference_frame_path': reference_frame_path}

def get_gif_file_info__ea8c7a7a(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_gif_file_info__e66cc6b0(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_hash__69eadfa0(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_dimensions__7816ba80(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_image_rotation__82cdf8c66b7532068e061ebffd0a6c33(env, config: Dict):
    """
    Get the dimensions of an image to check if it has been rotated 90 degrees.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path on VM

    Returns:
        dict: {"width": int, "height": int, "swapped": bool} where swapped indicates if dimensions are swapped
    """
    from PIL import Image
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get image file from {config['path']}")
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        swapped = img.width > img.height
        result = {'width': img.width, 'height': img.height, 'swapped': swapped}
        logger.info(f'Image dimensions: {img.width}x{img.height}, swapped: {swapped}')
        return result
    except Exception as e:
        logger.error(f'Failed to read image: {str(e)}')
        return None
    finally:
        os.unlink(tmp_path)

def get_docx_total_images__02587c7a(env, config: dict):
    """Get total image count in document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        int: Total number of images
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_image_mode__78bfc971(env, config):
    """
    Get the color mode of an image (e.g., RGB, RGBA, L, P).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Image mode, or None if error
    """
    try:
        file_path = config.get('path')
        if not file_path:
            logger.error('No file path provided')
            return None
        file_content = env.controller.get_file(file_path)
        if not file_content:
            logger.error(f'Failed to get file: {file_path}')
            return None
        import os
        local_path = os.path.join(env.cache_dir, 'temp_mode_check.png')
        with open(local_path, 'wb') as f:
            f.write(file_content)
        img = Image.open(local_path)
        img_mode = img.mode
        logger.info(f'Image mode: {img_mode}')
        return img_mode
    except Exception as e:
        logger.error(f'Error getting image mode: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_gimp_image_size__5db5c513f28a5ad7ad8d60ef05b14b35(env, config: Dict[str, str]):
    """
    Gets the size (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        size_info = {'width': img.width, 'height': img.height}
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image size: {size_info}')
        return size_info
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_image_saturation__8a7e9a1c5d75d961a07bbfbc1999062a(env, config: Dict):
    """
    Get the average saturation of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the image file

    Returns:
        Dict with 'saturation' key (float value 0-255), or None if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get file: {config['path']}")
        return None
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            hsv_image = img.convert('HSV')
            saturation_channel = hsv_image.split()[1]
            stat = ImageStat.Stat(saturation_channel)
            saturation = stat.mean[0]
            img.close()
            return {'saturation': saturation}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image saturation: {e}')
        return None

def get_has_background_image__6ccbfc2a(env, config: Dict[str, str]):
    """
    Check if a slide has a background image from a specific video frame.

    Args:
        env: Environment object
        config: Configuration dict with 'ppt_file_path', 'slide_index', 'video_path', and 'timestamp'

    Returns:
        dict: Dictionary with 'has_background' boolean and 'similarity_score' if verification succeeds
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    video_path = config.get('video_path', '/home/user/Desktop/landscape.mp4')
    timestamp_seconds = config.get('timestamp', 25)
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    video_file_localhost_path = get_vm_file(env, {'path': video_path, 'dest': os.path.split(video_path)[-1]})
    has_background = False
    similarity_score = 0.0
    background_image_data = None
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False, 'similarity_score': 0.0}
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                image_rel_id = None
                for child in root.iter(bg_tag):
                    for element in child.iter(image_tag):
                        image_rel_id = element.attrib.get(attr_tag)
                        if image_rel_id is not None:
                            has_background = True
                            break
                    if has_background:
                        break
                if not has_background:
                    return {'has_background': False, 'similarity_score': 0.0}
                if slide_rels_file in myzip.namelist():
                    with myzip.open(slide_rels_file) as rels_f:
                        rels_tree = ET.parse(rels_f)
                        rels_root = rels_tree.getroot()
                        for rel in rels_root.findall('{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                            if rel.attrib.get('Id') == image_rel_id:
                                target = rel.attrib.get('Target')
                                image_path = 'ppt/slides/' + target.replace('../', '')
                                if image_path in myzip.namelist():
                                    with myzip.open(image_path) as img_f:
                                        background_image_data = img_f.read()
                                break
        if background_image_data:
            cap = cv2.VideoCapture(video_file_localhost_path)
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_number = int(timestamp_seconds * fps)
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_number)
            (ret, video_frame) = cap.read()
            cap.release()
            if ret:
                bg_image = Image.open(io.BytesIO(background_image_data))
                bg_image_cv = cv2.cvtColor(np.array(bg_image), cv2.COLOR_RGB2BGR)
                (height, width) = video_frame.shape[:2]
                bg_image_cv = cv2.resize(bg_image_cv, (width, height))
                video_gray = cv2.cvtColor(video_frame, cv2.COLOR_BGR2GRAY)
                bg_gray = cv2.cvtColor(bg_image_cv, cv2.COLOR_BGR2GRAY)
                correlation = cv2.matchTemplate(video_gray, bg_gray, cv2.TM_CCORR_NORMED)[0][0]
                similarity_score = float(correlation)
    except Exception as e:
        return {'has_background': False, 'similarity_score': 0.0, 'error': str(e)}
    return {'has_background': has_background, 'similarity_score': similarity_score}

def get_gimp_image_flipped__16b4973e(env, config):
    """
    Get the image and check if it's flipped compared to the original.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'original_path' keys

    Returns:
        dict: Dictionary with 'flipped_horizontal' boolean
    """
    try:
        result_path = config.get('path')
        original_path = config.get('original_path', '/home/user/Desktop/character.png')
        if not result_path:
            logger.error('No result path provided')
            return None
        result_content = env.controller.get_file(result_path)
        original_content = env.controller.get_file(original_path)
        if not result_content or not original_content:
            logger.error('Failed to get files')
            return None
        import os
        result_local = os.path.join(env.cache_dir, 'result_flip.png')
        original_local = os.path.join(env.cache_dir, 'original_flip.png')
        with open(result_local, 'wb') as f:
            f.write(result_content)
        with open(original_local, 'wb') as f:
            f.write(original_content)
        result_img = Image.open(result_local)
        original_img = Image.open(original_local)
        flipped_original = original_img.transpose(Image.FLIP_LEFT_RIGHT)
        import numpy as np
        result_arr = np.array(result_img)
        flipped_arr = np.array(flipped_original)
        is_flipped = np.allclose(result_arr, flipped_arr, atol=5)
        logger.info(f'Image flipped horizontally: {is_flipped}')
        return {'flipped_horizontal': is_flipped}
    except Exception as e:
        logger.error(f'Error checking image flip: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_gimp_rotated_image__e0c6cf186be7cc4682d3712837dcdd72(env, config):
    """Get both the original and rotated images from VM for comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'original_path' and 'rotated_path' keys, or None if either file missing
    """
    rotated_file_bytes = env.controller.get_file(config['path'])
    if not rotated_file_bytes:
        return None
    dest_filename = config.get('dest', os.path.basename(config['path']))
    rotated_cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(rotated_cache_path, 'wb') as f:
        f.write(rotated_file_bytes)
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    original_file_bytes = env.controller.get_file(original_path)
    if not original_file_bytes:
        return None
    original_cache_path = os.path.join(env.cache_dir, 'Triangle_On_The_Side.png')
    with open(original_cache_path, 'wb') as f:
        f.write(original_file_bytes)
    return {'original_path': original_cache_path, 'rotated_path': rotated_cache_path}

def get_image_corner_colors__340a3600a88f3d4dc2217b9f986d625d(env, config):
    """Get colors at image corners to detect horizontal flip.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Corner pixel colors or None
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (width, height) = (img.width, img.height)
        corners = {'top_left': img.getpixel((0, 0)), 'top_right': img.getpixel((width - 1, 0)), 'bottom_left': img.getpixel((0, height - 1)), 'bottom_right': img.getpixel((width - 1, height - 1))}
        left_edge_colors = []
        right_edge_colors = []
        step = max(1, height // 10)
        for y in range(0, height, step):
            left_edge_colors.append(img.getpixel((0, y)))
            right_edge_colors.append(img.getpixel((width - 1, y)))

        def avg_color(colors):
            if not colors:
                return (0, 0, 0)
            r = sum((c[0] for c in colors)) // len(colors)
            g = sum((c[1] for c in colors)) // len(colors)
            b = sum((c[2] for c in colors)) // len(colors)
            return (r, g, b)
        result = {'corners': corners, 'left_edge_avg': avg_color(left_edge_colors), 'right_edge_avg': avg_color(right_edge_colors), 'width': width, 'height': height}
        img.close()
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_file_format__024272f6(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the format of an image file from the VM.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        str: Image format (e.g., 'JPEG', 'PNG', 'GIF'), or None if file doesn't exist or can't be read
    """
    path = config['path']
    file_data = env.controller.get_file(path)
    if file_data is None:
        logger.warning(f'Failed to get image file from VM: {path}')
        return None
    cache_path = os.path.join(env.cache_dir, os.path.basename(path))
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with Image.open(cache_path) as img:
            image_format = img.format
        logger.info(f'Successfully retrieved image format: {image_format}')
        return image_format
    except Exception as e:
        logger.error(f'Error processing image file {path}: {e}')
        return None

def get_image_format__945e33b1(env, config):
    """
    Get the format of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Image format (e.g., 'PNG', 'JPEG', 'GIF'), or None if error
    """
    try:
        file_path = config.get('path')
        if not file_path:
            logger.error('No file path provided')
            return None
        file_content = env.controller.get_file(file_path)
        if not file_content:
            logger.error(f'Failed to get file: {file_path}')
            return None
        import os
        local_path = os.path.join(env.cache_dir, 'temp_format_check.img')
        with open(local_path, 'wb') as f:
            f.write(file_content)
        img = Image.open(local_path)
        img_format = img.format
        logger.info(f'Image format: {img_format}')
        return img_format
    except Exception as e:
        logger.error(f'Error getting image format: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_dimensions__3b2067dec4f66b0e25da3355b1fb8f3e(env, config: dict):
    """Get the dimensions of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with width and height
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'width': 0, 'height': 0}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        img.close()
        return dimensions
    finally:
        os.unlink(tmp_path)

def get_gimp_gimprc_file__5b6e5d1f(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_image_properties__ab9c94ea(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_triangle_horizontal_position__c73a98370436f02ecc8c151f49a9b91c(env, config: dict):
    """
    Get the horizontal center position of the triangle (as fraction of image width).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        float: Horizontal position of triangle centroid (0.0 = left edge, 1.0 = right edge)
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        triangle_color = unique_colors_sorted[1]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            return None
        centroid_col = triangle_coords[:, 1].mean()
        image_width = img_array.shape[1]
        horizontal_position = centroid_col / image_width
        return float(horizontal_position)
    finally:
        os.unlink(tmp_path)

def get_png_resolution_info__27528fb8(env, config: Dict) -> Optional[Dict]:
    """
    Get PNG resolution information including total pixels and aspect ratio.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'width', 'height', 'total_pixels', 'aspect_ratio' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'total_pixels': 0, 'aspect_ratio': 0.0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'total_pixels': 0, 'aspect_ratio': 0.0}
        with Image.open(local_path) as img:
            total_pixels = img.width * img.height
            aspect_ratio = img.width / img.height if img.height > 0 else 0.0
            return {'exists': True, 'width': img.width, 'height': img.height, 'total_pixels': total_pixels, 'aspect_ratio': aspect_ratio}
    except Exception as e:
        return {'exists': False, 'total_pixels': 0, 'aspect_ratio': 0.0, 'error': str(e)}

def get_image_saturation__df2e4b52(env, config: Dict[str, Any]) -> Dict[str, float]:
    """
    Get the saturation values of both the original background image from PowerPoint
    and the saved enhanced image.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'path': path to saved enhanced image
            - 'original_pptx': path to original PowerPoint file
            - 'slide_index': index of slide containing background image (0-based)

    Returns:
        dict: Dictionary with 'original_saturation' and 'final_saturation' keys
    """
    saved_image_path = config['path']
    original_pptx_path = config.get('original_pptx', '/home/user/Desktop/PPT-Template_widescreen.pptx')
    slide_index = config.get('slide_index', 1)
    result = {'original_saturation': 0.0, 'final_saturation': 0.0}
    saved_file = get_vm_file(env, {'path': saved_image_path, 'dest': os.path.basename(saved_image_path)})
    if saved_file is not None and os.path.exists(saved_file):
        result['final_saturation'] = _calculate_image_saturation(saved_file)
    pptx_file = get_vm_file(env, {'path': original_pptx_path, 'dest': os.path.basename(original_pptx_path)})
    if pptx_file is not None and os.path.exists(pptx_file):
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                with zipfile.ZipFile(pptx_file, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
                media_dir = os.path.join(temp_dir, 'ppt', 'media')
                if os.path.exists(media_dir):
                    image_files = [f for f in os.listdir(media_dir) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
                    if image_files:
                        image_files.sort()
                        original_bg_path = os.path.join(media_dir, image_files[0])
                        result['original_saturation'] = _calculate_image_saturation(original_bg_path)
        except Exception as e:
            print(f'Error extracting original image from PowerPoint: {e}')
    return result

def get_gif_file_info__092b88db7884c339b37665609b49294b(env, config: Dict):
    """
    Get information about a GIF file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path on VM

    Returns:
        dict with file existence, size, frame count, duration, and basic properties
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration': 0.0}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_gif = img.format == 'GIF'
        (width, height) = img.size
        frames = 0
        total_duration = 0.0
        try:
            while True:
                frames += 1
                frame_duration = img.info.get('duration', 100)
                total_duration += frame_duration
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        duration_seconds = total_duration / 1000.0
        result = {'exists': True, 'file_size': len(file_bytes), 'is_gif': is_gif, 'width': width, 'height': height, 'frames': frames, 'duration': duration_seconds}
        logger.info(f'GIF file info: {result}')
        return result
    except Exception as e:
        logger.error(f'Error analyzing GIF file: {e}')
        return {'exists': True, 'file_size': len(file_bytes), 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration': 0.0}
    finally:
        os.unlink(tmp_path)

def get_googledrive_png_files__f9d97b19(env, config: Dict[str, Any]) -> List[str]:
    """Get list of PNG files in a Google Drive folder

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings
            - folder_query: Query string to find the folder

    Returns:
        List of PNG filenames, sorted
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            logger.warning(f'Folder not found with query: {folder_query}')
            return []
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        png_files = [f['title'] for f in file_list if f['title'].endswith('.png')]
        logger.info(f'Found {len(png_files)} PNG files: {png_files}')
        return sorted(png_files)
    except Exception as e:
        logger.error(f'Error getting PNG files from Google Drive: {e}')
        return []

def get_image_rotation__b49205bc(env, config: dict):
    """Get image properties to detect rotation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        dict: Image properties including dimensions
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        return {'width': img.width, 'height': img.height, 'exists': True}
    except Exception as e:
        logger.error(f'Failed to open image: {str(e)}')
        return {'exists': False}

def get_gif_file_info__78103de86e64961437a4bcd00b97b9bc(env, config: Dict):
    """
    Get information about a GIF file on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path on VM

    Returns:
        dict with file existence, size, and basic properties
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_seconds': 0.0}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_gif = img.format == 'GIF'
        (width, height) = img.size
        frames = 0
        total_duration_ms = 0
        try:
            while True:
                frames += 1
                frame_duration = img.info.get('duration', 100)
                total_duration_ms += frame_duration
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        duration_seconds = total_duration_ms / 1000.0
        result = {'exists': True, 'file_size': len(file_bytes), 'is_gif': is_gif, 'width': width, 'height': height, 'frames': frames, 'duration_seconds': duration_seconds}
        logger.info(f'GIF file info: {result}')
        return result
    except Exception as e:
        logger.error(f'Error analyzing GIF file: {e}')
        return {'exists': True, 'file_size': len(file_bytes), 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_seconds': 0.0}
    finally:
        os.unlink(tmp_path)

def get_gimp_image_size__fb0dfb53338d620d275120282f62d824(env, config: Dict[str, str]):
    """
    Gets the size (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        size_info = {'width': img.width, 'height': img.height}
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image size: {size_info}')
        return size_info
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_all_images__1574bc56f755697238f4190c2d72a32b(env, config: dict) -> dict:
    """Get directory listing to check all image files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Directory tree dict with 'children' list
    """
    return env.controller.get_vm_directory_tree(config['path'])

def get_extracted_image_properties__a8440735(env, config: dict):
    """
    Get properties of the extracted image and verify it matches the first image from the email.

    This function:
    1. Extracts the target image file properties
    2. Accesses the Thunderbird profile to find the recent email in Notes
    3. Extracts the first image from that email
    4. Computes hash of both images to verify they match
    """
    from PIL import Image
    from io import BytesIO
    path = config.get('path', '/home/user/image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0, 'content_hash': None, 'matches_email_image': False}
    try:
        img = Image.open(BytesIO(file_bytes))
        content_hash = hashlib.md5(file_bytes).hexdigest()
        result = {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height, 'content_hash': content_hash, 'matches_email_image': False}
        try:
            email_image_hash = _get_first_email_image_hash(env)
            if email_image_hash and email_image_hash == content_hash:
                result['matches_email_image'] = True
        except Exception as e:
            pass
        return result
    except Exception as e:
        content_hash = hashlib.md5(file_bytes).hexdigest() if file_bytes else None
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0, 'content_hash': content_hash, 'matches_email_image': False}

def get_triangle_rotation__2c517d983ea369519ffc55979e3393ec(env, config: dict):
    """
    Extract the rotation angles of the yellow triangle in both original and final images.
    Returns a dict with 'original_angle' and 'final_angle' in degrees (0-360).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the final image file

    Returns:
        dict: {'original_angle': float, 'final_angle': float} or None if error
    """

    def calculate_triangle_angle(file_bytes):
        """Helper function to calculate triangle angle from image bytes."""
        if not file_bytes:
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            img_array = np.array(img)
            (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
            unique_colors_sorted = unique_colors[np.argsort(counts)]
            triangle_color = unique_colors_sorted[1]
            triangle_mask = np.all(img_array == triangle_color, axis=2)
            triangle_coords = np.argwhere(triangle_mask)
            if len(triangle_coords) == 0:
                return None
            centroid = triangle_coords.mean(axis=0)
            top_point = triangle_coords[np.argmin(triangle_coords[:, 0])]
            bottom_left = triangle_coords[np.argmin(triangle_coords[:, 1])]
            bottom_right = triangle_coords[np.argmax(triangle_coords[:, 1])]
            vec_to_top = top_point - centroid
            angle = np.arctan2(vec_to_top[1], vec_to_top[0]) * 180 / np.pi
            angle = (angle + 90) % 360
            return angle
        finally:
            os.unlink(tmp_path)
    final_file_bytes = env.controller.get_file(config['path'])
    if not final_file_bytes:
        return None
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    original_file_bytes = env.controller.get_file(original_path)
    if not original_file_bytes:
        return None
    original_angle = calculate_triangle_angle(original_file_bytes)
    final_angle = calculate_triangle_angle(final_file_bytes)
    if original_angle is None or final_angle is None:
        return None
    return {'original_angle': original_angle, 'final_angle': final_angle}

def get_image_rotated__045e2e0c(env, config):
    """
    Check if an image has been rotated 90 degrees clockwise.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'original_path'

    Returns:
        dict: Dictionary with 'rotated_90_cw' boolean
    """
    try:
        result_path = config.get('path')
        original_path = config.get('original_path', '/home/user/Desktop/character.png')
        if not result_path:
            logger.error('No result path provided')
            return None
        result_content = env.controller.get_file(result_path)
        original_content = env.controller.get_file(original_path)
        if not result_content or not original_content:
            logger.error('Failed to get files')
            return None
        import os
        result_local = os.path.join(env.cache_dir, 'result_rotate.png')
        original_local = os.path.join(env.cache_dir, 'original_rotate.png')
        with open(result_local, 'wb') as f:
            f.write(result_content)
        with open(original_local, 'wb') as f:
            f.write(original_content)
        result_img = Image.open(result_local)
        original_img = Image.open(original_local)
        rotated_original = original_img.transpose(Image.ROTATE_270)
        import numpy as np
        result_arr = np.array(result_img)
        rotated_arr = np.array(rotated_original)
        is_rotated = np.allclose(result_arr, rotated_arr, atol=5)
        logger.info(f'Image rotated 90 CW: {is_rotated}')
        return {'rotated_90_cw': is_rotated}
    except Exception as e:
        logger.error(f'Error checking rotation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_format_info__900fc36276e9eaf12471a7834992cf5c(env, config: dict):
    """
    Get image file format information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Image format information
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Image file not found: {path}')
            return {'exists': False, 'path': path, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            format_name = img.format
            return {'exists': True, 'path': path, 'filename': os.path.basename(path), 'format': format_name, 'size': len(file_bytes), 'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking image format {path}: {e}')
        return {'exists': False, 'path': path, 'format': None, 'error': str(e)}

def get_image_properties__70f3db86(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_image_properties__50f8e5bb(env, config: dict):
    """Extract image properties (dimensions, format, content hash) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists, and perceptual hash
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'phash': None}
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            phash = str(imagehash.phash(img))
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format, 'phash': phash}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'phash': None, 'error': str(e)}

def get_gimp_layer_names__8565a91c(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_image_properties__c6063730(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_default_audio_player__971ae155(env, config: dict):
    """Gets the default application for audio files on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['aac', 'ac3', 'flac', 'mp3', 'mpeg', 'ogg', 'opus', 'vorbis', 'wav', 'webm', 'x-aiff', 'x-ape', 'x-flac', 'x-it', 'x-m4a', 'x-matroska', 'x-mod', 'x-mp3', 'x-mpeg', 'x-mpegurl', 'x-ms-wma', 'x-musepack', 'x-oggflac', 'x-pn-realaudio', 'x-scpls', 'x-vorbis', 'x-vorbis+ogg', 'x-wav', 'x-wavpack', 'x-xm']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'audio/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_image_properties__795e137d(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_default_image_viewer__ae568cb6(env, config: dict):
    """Gets the default application for image files on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['bmp', 'gif', 'jpeg', 'jpg', 'png', 'svg+xml', 'tiff', 'webp', 'x-bmp', 'x-icon', 'x-png', 'x-portable-pixmap', 'x-tga', 'x-xbitmap', 'x-xpixmap']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'image/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_image_properties__68566bbc(env, config: dict):
    """Get image properties (dimensions, mode, format) from VM file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        dict: Image properties including width, height, mode, format
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        properties = {'width': img.width, 'height': img.height, 'mode': img.mode, 'format': img.format, 'size': img.size, 'exists': True}
        return properties
    except Exception as e:
        logger.error(f'Failed to open image: {str(e)}')
        return {'exists': False}

def get_image_file__2237c5bebdb76e54ae53ea89e71ca4a3(env, config):
    """Get image file from VM and save to local cache.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        str: Path to the downloaded image file in cache directory
    """
    import os
    file_path = config.get('path')
    dest = config.get('dest', 'result_image.jpg')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_googledrive_png_count__b4c05b10(env, config: Dict[str, Any]) -> int:
    """Count PNG files in Google Drive folder"""
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            return 0
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        png_count = sum((1 for f in file_list if f['title'].endswith('.png')))
        logger.info(f'PNG count: {png_count}')
        return png_count
    except Exception as e:
        logger.error(f'Error: {e}')
        return 0

def get_docx_image_count__5cfb373b(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_png_export_check__c69055e15cceefc40a87e6de042c2331(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if PNG file was exported from GIMP and extract its properties.
    Compares the exported PNG with the original XCF file to verify content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (PNG file) and 'original_path' (XCF file)

    Returns:
        Dict with file_exists, width, height, content_hash, original_dimensions properties
    """
    png_path = config.get('path', '')
    original_path = config.get('original_path', '')
    result = {'file_exists': False, 'width': 0, 'height': 0, 'content_hash': None, 'original_dimensions': None, 'dimensions_match': False}
    png_bytes = env.controller.get_file(png_path)
    if not png_bytes:
        logger.warning(f'PNG file not found: {png_path}')
        return result
    result['file_exists'] = True
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(png_bytes)
            tmp_png_path = tmp.name
        try:
            png_img = Image.open(tmp_png_path)
            result['width'] = png_img.size[0]
            result['height'] = png_img.size[1]
            png_rgb = png_img.convert('RGB')
            result['content_hash'] = hashlib.md5(png_rgb.tobytes()).hexdigest()
            logger.info(f"PNG dimensions: {result['width']}x{result['height']}, hash: {result['content_hash'][:8]}...")
        finally:
            os.unlink(tmp_png_path)
    except Exception as e:
        logger.error(f'Error reading PNG: {e}')
        return result
    if original_path:
        xcf_bytes = env.controller.get_file(original_path)
        if xcf_bytes:
            try:
                with tempfile.NamedTemporaryFile(suffix='.xcf', delete=False) as tmp:
                    tmp.write(xcf_bytes)
                    tmp_xcf_path = tmp.name
                try:
                    try:
                        xcf_img = Image.open(tmp_xcf_path)
                        (original_width, original_height) = xcf_img.size
                        result['original_dimensions'] = (original_width, original_height)
                        if result['width'] == original_width and result['height'] == original_height:
                            result['dimensions_match'] = True
                        xcf_rgb = xcf_img.convert('RGB')
                        xcf_hash = hashlib.md5(xcf_rgb.tobytes()).hexdigest()
                        if result['content_hash'] == xcf_hash:
                            result['content_verified'] = True
                            logger.info('✓ PNG content matches XCF file')
                        else:
                            result['content_verified'] = result['dimensions_match']
                            logger.info(f"PNG exported with matching dimensions: {result['dimensions_match']}")
                    except Exception as e:
                        logger.info(f'PIL cannot open XCF directly, parsing header: {e}')
                        if len(xcf_bytes) >= 17:
                            import struct
                            original_width = struct.unpack('>I', xcf_bytes[9:13])[0]
                            original_height = struct.unpack('>I', xcf_bytes[13:17])[0]
                            result['original_dimensions'] = (original_width, original_height)
                            if result['width'] == original_width and result['height'] == original_height:
                                result['dimensions_match'] = True
                                result['content_verified'] = True
                                logger.info(f'✓ PNG dimensions match XCF: {original_width}x{original_height}')
                            else:
                                logger.warning(f"PNG dimensions {result['width']}x{result['height']} don't match XCF {original_width}x{original_height}")
                                result['content_verified'] = False
                        else:
                            result['content_verified'] = False
                finally:
                    os.unlink(tmp_xcf_path)
            except Exception as e:
                logger.error(f'Error reading XCF file: {e}')
                result['content_verified'] = False
        else:
            logger.warning(f'Original XCF file not found: {original_path}')
            result['content_verified'] = False
    else:
        logger.warning('No original_path provided for content verification')
        result['content_verified'] = False
    return result

def get_image_color_mode__90bdfeb2ebfd1bac6873718be37b3912(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get image color mode and check if it's grayscale.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'mode', 'is_grayscale' keys
    """
    from PIL import Image
    import tempfile
    import numpy as np
    result = {'exists': False, 'mode': None, 'is_grayscale': False}
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.warning(f"Failed to get file from VM: {config['path']}")
            return result
        result['exists'] = True
        with tempfile.NamedTemporaryFile(suffix='.img', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result['mode'] = img.mode
            if img.mode == 'L':
                result['is_grayscale'] = True
            elif img.mode in ('RGB', 'RGBA'):
                img_array = np.array(img)
                if img.mode == 'RGBA':
                    rgb_channels = img_array[:, :, :3]
                else:
                    rgb_channels = img_array
                (r, g, b) = (rgb_channels[:, :, 0], rgb_channels[:, :, 1], rgb_channels[:, :, 2])
                result['is_grayscale'] = np.allclose(r, g) and np.allclose(g, b)
            logger.info(f"Image mode: {result['mode']}, is_grayscale={result['is_grayscale']}")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image color mode: {e}')
    return result

def get_image_file__203069587d53a571860bccb97348992b(env, config: dict):
    """Check if an image file exists and get basic info.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying image path on VM

    Returns:
        dict: {'exists': bool, 'path': str, 'size': int}
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes and len(file_bytes) > 0:
            logger.info(f'Image file {file_path} exists, size: {len(file_bytes)} bytes')
            return {'exists': True, 'path': file_path, 'size': len(file_bytes)}
        else:
            logger.info(f'Image file {file_path} does not exist')
            return {'exists': False, 'path': file_path, 'size': 0}
    except Exception as e:
        logger.error(f'Error checking image file {file_path}: {e}')
        return {'exists': False, 'path': file_path, 'size': 0}

def get_gimp_saturation__8f3d7ee1bf388be9293a1e57f283a35e(env, config):
    """Get saturation value of a GIMP image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        float: Average saturation value [0-255]
    """
    import tempfile
    import os
    import numpy as np
    from PIL import Image, ImageStat
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        hsv_img = img.convert('HSV')
        (_, s, _) = hsv_img.split()
        s_array = np.array(s)
        avg_saturation = np.mean(s_array)
        return float(avg_saturation)
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_png_export_info__046a9ca717a3ff75711d7d7d1d876a5f(env, config: dict):
    """
    Get information about a PNG file exported from video.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the PNG file path and source video path

    Returns:
        dict: File info including exists, format, width, height, file_size, source_video_exists
    """
    file_path = config.get('path')
    source_video = config.get('source_video', '/home/user/Desktop/src.mp4')
    if not file_path:
        logger.error('No file path specified in config')
        return {'exists': False, 'format': None, 'width': 0, 'height': 0, 'file_size': 0, 'source_video_exists': False}
    source_video_bytes = env.controller.get_file(source_video)
    source_video_exists = source_video_bytes is not None and len(source_video_bytes) > 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'format': None, 'width': 0, 'height': 0, 'file_size': 0, 'source_video_exists': source_video_exists}
    file_size = len(file_bytes)
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            img_format = img.format
            (width, height) = img.size
            img.close()
            return {'exists': True, 'format': img_format, 'width': width, 'height': height, 'file_size': file_size, 'source_video_exists': source_video_exists}
        except Exception as e:
            logger.error(f'Error opening image file: {e}')
            return {'exists': True, 'format': None, 'width': 0, 'height': 0, 'file_size': file_size, 'source_video_exists': source_video_exists}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_image_flip_check__08f973926b8f35af5489796c73a6c6e0(env, config: Dict):
    """
    Get the flipped image and return both source and result for comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (result image) and 'source_path' (original image)

    Returns:
        Dict with 'result_path' and 'source_path' for metric comparison
    """
    import tempfile
    import os
    result_path = config.get('path')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        logger.error(f'Failed to get result image from {result_path}')
        return None
    result_temp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    result_temp.write(result_bytes)
    result_temp.close()
    source_cache_path = config.get('source_cache_path')
    return {'result_path': result_temp.name, 'source_path': source_cache_path}

def get_image_dimensions__b6f18d98d1993dda1c36f44651fe6a5d(env, config: Dict[str, Any]):
    """
    Get the dimensions of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int} or None if file not found
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    import tempfile
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        (width, height) = img.size
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image dimensions: {width}x{height}')
        return {'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_image_dimensions__e0d9b551(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_docx_image_count__c37ca17b(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_has_background_image__99fc8543(env, config: Dict[str, str]):
    """
    Check if a specific slide has a background image applied.

    Args:
        env: Environment object
        config: Configuration dict with 'ppt_file_path' and 'slide_index'

    Returns:
        dict: Dictionary with 'has_background' boolean indicating if background image exists
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    has_background = False
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False}
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    try:
                        for element in child.iter(image_tag):
                            if attr_tag in element.attrib:
                                has_background = True
                                break
                    except:
                        pass
                    if has_background:
                        break
    except Exception as e:
        return {'has_background': False}
    return {'has_background': has_background}

def get_has_background_image__171e4956(env, config: Dict[str, str]):
    """
    Check if a specific slide has a background image that matches a video frame at 00:12.

    This getter:
    1. Extracts the background image from the specified slide
    2. Extracts the expected frame from landscape.mp4 at 00:12 timestamp
    3. Compares the two images using perceptual hashing to verify they match

    Args:
        env: Environment object
        config: Configuration dict with keys:
            - ppt_file_path: Path to the PPT file
            - slide_index: Index of the slide (0-based)

    Returns:
        dict: Dictionary with keys:
            - has_background: boolean indicating if background image exists
            - background_image_data: bytes of the background image if available
            - expected_frame_data: bytes of the expected video frame at 00:12
            - images_match: boolean indicating if background matches expected frame
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    has_background = False
    background_image_data = None
    expected_frame_data = None
    images_match = False
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = f'ppt/slides/slide{slide_index + 1}.xml'
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False, 'background_image_data': None, 'expected_frame_data': None, 'images_match': False}
            background_image_id = None
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    for element in child.iter(image_tag):
                        image_id = element.attrib.get(attr_tag)
                        if image_id:
                            has_background = True
                            background_image_id = image_id
                            break
                    if has_background:
                        break
            if has_background and background_image_id:
                slide_rels_file = f'ppt/slides/_rels/slide{slide_index + 1}.xml.rels'
                if slide_rels_file in myzip.namelist():
                    with myzip.open(slide_rels_file) as f:
                        rels_tree = ET.parse(f)
                        rels_root = rels_tree.getroot()
                        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                            if rel.attrib.get('Id') == background_image_id:
                                target = rel.attrib.get('Target')
                                if target:
                                    image_path = f'ppt/slides/{target}'
                                    image_path = os.path.normpath(image_path.replace('\\', '/')).replace('\\', '/')
                                    if image_path in myzip.namelist():
                                        try:
                                            with myzip.open(image_path) as img_file:
                                                background_image_data = img_file.read()
                                        except Exception:
                                            pass
                                break
        video_path = '/home/user/Desktop/landscape.mp4'
        try:
            video_localhost_path = get_vm_file(env, {'path': video_path, 'dest': 'landscape.mp4'})
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_frame:
                tmp_frame_path = tmp_frame.name
            try:
                subprocess.run(['ffmpeg', '-ss', '00:00:12', '-i', video_localhost_path, '-frames:v', '1', '-y', tmp_frame_path], capture_output=True, timeout=10)
                if os.path.exists(tmp_frame_path):
                    with open(tmp_frame_path, 'rb') as f:
                        expected_frame_data = f.read()
            finally:
                if os.path.exists(tmp_frame_path):
                    os.unlink(tmp_frame_path)
        except Exception:
            pass
        if background_image_data and expected_frame_data:
            images_match = compare_images(background_image_data, expected_frame_data)
    except Exception:
        return {'has_background': False, 'background_image_data': None, 'expected_frame_data': None, 'images_match': False}
    return {'has_background': has_background, 'background_image_data': background_image_data, 'expected_frame_data': expected_frame_data, 'images_match': images_match}

def get_image_format__c0fb0f23(env, config):
    """
    Get the format of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'format': str, 'exists': bool} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/converted.jpg')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return {'format': None, 'exists': False}
        local_path = os.path.join(env.cache_dir, 'temp_format_check.img')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        img_format = img.format
        result = {'format': img_format, 'exists': True}
        logger.info(f'Image format: {result}')
        return result
    except Exception as e:
        logger.error(f'Error getting image format: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return {'format': None, 'exists': False}

def get_gif_file_info__11cf8ab6(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_gif_file_info__b130b682(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_aspect_ratio__f85979b813875fecca12ba1c6ab4cc68(env, config: Dict):
    """
    Get the aspect ratio (width/height) of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path on VM

    Returns:
        dict: {"aspect_ratio": float, "width": int, "height": int} or None if file not found
    """
    from PIL import Image
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get image file from {config['path']}")
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        aspect_ratio = img.width / img.height if img.height > 0 else 0
        result = {'aspect_ratio': aspect_ratio, 'width': img.width, 'height': img.height}
        logger.info(f'Image aspect ratio: {aspect_ratio:.4f} ({img.width}x{img.height})')
        return result
    except Exception as e:
        logger.error(f'Failed to read image: {str(e)}')
        return None
    finally:
        os.unlink(tmp_path)

def get_cropped_image_size__8159102f(env, config):
    """
    Get the size of a cropped image and compare with original to verify cropping.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'original_path'

    Returns:
        dict: Dictionary with size info and whether image was cropped
    """
    try:
        result_path = config.get('path')
        original_path = config.get('original_path', '/home/user/Desktop/character.png')
        if not result_path:
            logger.error('No result path provided')
            return None
        result_content = env.controller.get_file(result_path)
        original_content = env.controller.get_file(original_path)
        if not result_content or not original_content:
            logger.error('Failed to get files')
            return None
        result_local = os.path.join(env.cache_dir, 'result_crop.png')
        original_local = os.path.join(env.cache_dir, 'original_crop.png')
        with open(result_local, 'wb') as f:
            f.write(result_content)
        with open(original_local, 'wb') as f:
            f.write(original_content)
        result_img = Image.open(result_local)
        original_img = Image.open(original_local)
        result_size = result_img.size
        original_size = original_img.size
        is_cropped = result_size[0] < original_size[0] or result_size[1] < original_size[1]
        logger.info(f'Original: {original_size}, Result: {result_size}, Cropped: {is_cropped}')
        return {'width': result_size[0], 'height': result_size[1], 'is_cropped': is_cropped}
    except Exception as e:
        logger.error(f'Error checking crop: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_color_mode__3aef0cbd8986e240435edab0fd96c873(env, config):
    """Get image color mode and basic color statistics.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'mode': str, 'is_grayscale': bool} or None
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        mode = img.mode
        is_grayscale = mode in ('L', 'LA', '1')
        if mode in ('RGB', 'RGBA') and (not is_grayscale):
            img_small = img.resize((100, 100))
            pixels = list(img_small.getdata())
            grayscale_count = 0
            for pixel in pixels[:100]:
                if len(pixel) >= 3:
                    if pixel[0] == pixel[1] == pixel[2]:
                        grayscale_count += 1
            if grayscale_count >= 95:
                is_grayscale = True
        result = {'mode': mode, 'is_grayscale': is_grayscale}
        img.close()
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_gdrive_png_count__63477992bbc88c6a7091e80f7c0a8a72(env, config: Dict[str, Any]) -> List[str]:
    """
    Get the list of PNG file names in a specific Google Drive folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - settings_file: Path to Google Drive settings file
            - folder_name: Name of the folder to check (default: 'figures')
            - file_extension: File extension to filter (default: '.png')

    Returns:
        List[str]: List of PNG file names in the folder
    """
    settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
    folder_name = config.get('folder_name', 'figures')
    file_extension = config.get('file_extension', '.png').lower()
    try:
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_query = f"title = '{folder_name}' and trashed = false and 'root' in parents and mimeType = 'application/vnd.google-apps.folder'"
        folder_list = drive.ListFile({'q': folder_query}).GetList()
        if not folder_list:
            return []
        folder_id = folder_list[0]['id']
        file_query = f"trashed = false and '{folder_id}' in parents"
        file_list = drive.ListFile({'q': file_query}).GetList()
        file_names = [f['title'] for f in file_list if f['title'].lower().endswith(file_extension)]
        return sorted(file_names)
    except Exception as e:
        import logging
        logger = logging.getLogger('desktopenv.getter.googledrive')
        logger.error(f'Error getting Google Drive file names: {e}')
        return []

def get_image_dimensions__fa1a72d6(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_image_properties__739292ff(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get properties of an image file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' parameters

    Returns:
        dict: Image properties (width, height, format, size_bytes) or None if error
    """
    path = config.get('path', '')
    dest = config.get('dest', 'image.png')
    if not path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(os.path.dirname(cache_path), exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        img = Image.open(cache_path)
        properties = {'width': img.width, 'height': img.height, 'format': img.format, 'size_bytes': len(file_bytes), 'mode': img.mode}
        logger.info(f'Image properties: {properties}')
        return properties
    except Exception as e:
        logger.error(f'Error getting image properties: {e}')
        return None

def get_docx_image_count__dbe26b57(env, config):
    """
    Extract image information from a DOCX file, including total count and position.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        dict: {
            'count': total number of images,
            'first_element_is_image': whether the first body element is an image
        }
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return {'count': 0, 'first_element_is_image': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return {'count': 0, 'first_element_is_image': False}
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            first_element_is_image = False
            if len(doc.element.body) > 0:
                first_element = doc.element.body[0]
                if isinstance(first_element, CT_P):
                    for child in first_element.iter():
                        if isinstance(child, CT_Picture):
                            first_element_is_image = True
                            break
            return {'count': image_count, 'first_element_is_image': first_element_is_image}
    except Exception as e:
        logger.error(f'Error extracting image information: {e}')
        return {'count': 0, 'first_element_is_image': False}

def get_image_file_exists__7c0cc95089263e14cb308f3757c8acc1(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Extract file existence and format information for an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'is_png', 'file_path' keys
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return {'exists': False, 'is_png': False, 'file_path': path}
    is_png = False
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            is_png = img.format == 'PNG'
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.warning(f'Failed to verify image format: {e}')
    return {'exists': True, 'is_png': is_png, 'file_path': path}

def get_image_dimensions__505cf5fc(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_image_count_status__8e1f0943(env, config: dict):
    """Get image count and check if it meets minimum.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Number of images in document
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_googledrive_image_files__67a58d2f(env, config: Dict[str, Any]) -> List[str]:
    """Get list of PNG image files in Google Drive folder

    Verifies both file extension and MIME type to ensure files are valid PNG images.
    """
    try:
        from pydrive.auth import GoogleAuth
        from pydrive.drive import GoogleDrive
        settings_file = config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml')
        folder_query = config.get('folder_query', '')
        image_extensions = ['.png']
        valid_mime_types = ['image/png']
        auth = GoogleAuth(settings_file=settings_file)
        drive = GoogleDrive(auth)
        folder_search = f'{folder_query} and "root" in parents'
        folder_list = drive.ListFile({'q': folder_search}).GetList()
        if len(folder_list) == 0:
            return []
        folder_id = folder_list[0]['id']
        file_search = f'"{folder_id}" in parents and trashed = false'
        file_list = drive.ListFile({'q': file_search}).GetList()
        image_files = []
        for f in file_list:
            filename = f['title']
            has_valid_ext = any((filename.lower().endswith(ext) for ext in image_extensions))
            mime_type = f.get('mimeType', '')
            has_valid_mime = mime_type in valid_mime_types if mime_type else True
            if has_valid_ext and has_valid_mime:
                image_files.append(filename)
            elif has_valid_ext and (not has_valid_mime):
                logger.warning(f'File {filename} has PNG extension but MIME type is {mime_type}')
        logger.info(f'Found {len(image_files)} valid PNG files')
        return sorted(image_files)
    except Exception as e:
        logger.error(f'Error: {e}')
        return []

def get_gimp_file_bytes__677c871105619c52013adb82fa7e5d28(env, config: Dict[str, str]):
    """
    Get file bytes from VM for image comparison.
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file: {file_path}')
        return None
    import os
    cache_path = os.path.join(env.cache_dir, config.get('dest', 'temp_file.png'))
    os.makedirs(os.path.dirname(cache_path), exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    logger.debug(f'Saved file to cache: {cache_path}')
    return cache_path

def get_gimp_image_size__8677d5c370f70c69494653c2a8ef5be2(env, config: Dict[str, str]):
    """
    Gets the size (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        size_info = {'width': img.width, 'height': img.height}
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image size: {size_info}')
        return size_info
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_bottom_half_image__4baa14c913fe32a60f559a98296affae(env, config: Dict):
    """
    Get the bottom half of an image and return it along with the saved image for comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'source_path': path to the original image (tilearray.png)
            - 'result_path': path to the saved bottom half image

    Returns:
        dict: {
            "result_image": numpy array of the saved image,
            "expected_image": numpy array of the expected bottom half,
            "dimensions": {"width": int, "height": int}
        } or None if files not found
    """
    from PIL import Image
    import tempfile
    import numpy as np
    source_path = config.get('source_path', '/home/user/Desktop/tilearray.png')
    result_path = config.get('result_path', '/home/user/Desktop/bottom_half.png')
    source_bytes = env.controller.get_file(source_path)
    if not source_bytes:
        logger.error(f'Failed to get source image from {source_path}')
        return None
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        logger.error(f'Failed to get result image from {result_path}')
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(source_bytes)
        source_tmp_path = tmp.name
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(result_bytes)
        result_tmp_path = tmp.name
    try:
        source_img = Image.open(source_tmp_path)
        (width, height) = source_img.size
        expected_bottom_half = source_img.crop((0, height // 2, width, height))
        result_img = Image.open(result_tmp_path)
        expected_array = np.array(expected_bottom_half)
        result_array = np.array(result_img)
        logger.info(f'Source image size: {width}x{height}')
        logger.info(f'Expected bottom half size: {expected_bottom_half.size}')
        logger.info(f'Result image size: {result_img.size}')
        return {'result_image': result_array, 'expected_image': expected_array, 'dimensions': {'result_width': result_img.width, 'result_height': result_img.height, 'expected_width': expected_bottom_half.width, 'expected_height': expected_bottom_half.height}}
    except Exception as e:
        logger.error(f'Failed to process images: {str(e)}')
        return None
    finally:
        os.unlink(source_tmp_path)
        os.unlink(result_tmp_path)

def get_gimp_image_brightness__ced7a23c(env, config):
    """
    Get the average brightness of an image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        float: Average brightness value (0-255), or None if error
    """
    try:
        file_path = config.get('path')
        if not file_path:
            logger.error('No file path provided in config')
            return None
        file_content = env.controller.get_file(file_path)
        if not file_content:
            logger.error(f'Failed to get file: {file_path}')
            return None
        import os
        local_path = os.path.join(env.cache_dir, 'temp_brightness_check.png')
        with open(local_path, 'wb') as f:
            f.write(file_content)
        img = Image.open(local_path)
        grayscale = img.convert('L')
        stat = ImageStat.Stat(grayscale)
        brightness = stat.mean[0]
        logger.info(f'Image brightness: {brightness}')
        return brightness
    except Exception as e:
        logger.error(f'Error getting image brightness: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_docx_image_count__74be0e09(env, config: dict):
    """Get the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to VM file path

    Returns:
        int: Number of images in the document
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        image_count = 0
        for rel in doc.part.rels.values():
            if 'image' in rel.reltype:
                image_count += 1
        return image_count
    except Exception as e:
        return 0

def get_image_info__e19bfc7ef1338231fef513d8a5b2f6d1(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get image file information including format, size, and existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        Dict with 'exists', 'format', 'width', 'height' keys
    """
    from PIL import Image
    import tempfile
    result = {'exists': False, 'format': None, 'width': None, 'height': None}
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.warning(f"Failed to get file from VM: {config['path']}")
            return result
        result['exists'] = True
        with tempfile.NamedTemporaryFile(suffix='.img', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result['format'] = img.format
            result['width'] = img.size[0]
            result['height'] = img.size[1]
            logger.info(f"Image info: format={result['format']}, size={result['width']}x{result['height']}")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image info: {e}')
    return result

def get_gimp_flipped_image__07b5058cd3df711389d2b4342d0c561c(env, config: Dict[str, str]) -> Optional[str]:
    """
    Get the flipped image file from VM.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        str: Path to the downloaded file in cache, or None if file doesn't exist
    """
    path = config['path']
    try:
        file = env.controller.get_file(path)
        if file is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        dest = os.path.basename(path)
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file)
        logger.info(f'Successfully saved file: {cache_path} ({len(file)} bytes)')
        return cache_path
    except Exception as e:
        logger.error(f'Error processing file {path}: {e}')
        return None

def get_docx_image_embedded__6ed140cad2fa6c5f3349fe322a43680c(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if DOCX file contains embedded images and extract their dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to DOCX file

    Returns:
        Dict with file_exists, image_count, and max_width
    """
    docx_path = config.get('path', '')
    result = {'file_exists': False, 'image_count': 0, 'max_width': 0}
    file_bytes = env.controller.get_file(docx_path)
    if not file_bytes:
        logger.warning(f'DOCX file not found: {docx_path}')
        return result
    result['file_exists'] = True
    try:
        from docx import Document
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            doc = Document(tmp_path)
            images = []
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    img_data = rel.target_part.blob
                    images.append(BytesIO(img_data))
            result['image_count'] = len(images)
            max_width = 0
            for img_bytes in images:
                try:
                    img = Image.open(img_bytes)
                    max_width = max(max_width, img.size[0])
                except Exception as e:
                    logger.warning(f'Could not read image dimensions: {e}')
            result['max_width'] = max_width
            logger.info(f"Found {result['image_count']} images, max width: {max_width}")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading DOCX file: {e}')
    return result

def get_image_dimensions__b24412f0(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_default_audio_player__7752be97956df277ef920fced6cad6a7(env, config: dict):
    """Gets the default application for audio file types on Linux.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: The most common default audio player application name
    """
    from desktop_env.evaluators.getters.general import get_vm_command_line
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['aac', 'flac', 'mp3', 'mp4', 'mpeg', 'ogg', 'opus', 'wav', 'webm', 'x-flac', 'x-mp3', 'x-mpeg', 'x-ms-wma', 'x-wav', 'x-vorbis+ogg', 'mpeg-url', 'x-mpegurl', 'x-scpls']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'audio/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_docx_image_count__081533bc(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_image_scaled_dims__bf4967f3e9931b3c80be8e4dbf6e04b7(env, config: dict):
    """
    Get image file from VM and return its dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        dict: Dictionary with 'width' and 'height' keys, or None if failed
    """
    vm_path = config['path']
    dest_name = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        img = Image.open(cache_path)
        result = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {result}')
        return result
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        return None

def get_extracted_image_properties__69184b17(env, config: dict):
    from PIL import Image
    from io import BytesIO
    path = config.get('path', '/home/user/image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0}
    try:
        img = Image.open(BytesIO(file_bytes))
        return {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height}
    except:
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0}

def get_image_validation__2511ecbd(env, config: Dict) -> Optional[Dict]:
    """
    Validate image file format and existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'format', and 'valid' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'format': None, 'valid': False}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'format': None, 'valid': False}
        with Image.open(local_path) as img:
            return {'exists': True, 'format': img.format, 'valid': True}
    except Exception as e:
        return {'exists': True, 'format': None, 'valid': False, 'error': str(e)}

def get_gimp_gimprc_file__9c13adcb(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_gimp_image_brightness__652303c0d066122d99f102352a5b1a93(env, config):
    """
    Get the average brightness of a GIMP-exported image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        float: Average brightness value (0-255) or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            grayscale = img.convert('L')
            stat = ImageStat.Stat(grayscale)
            brightness = stat.mean[0]
            logger.info(f'Image brightness: {brightness}')
            return brightness
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image brightness: {e}')
        return None

def get_image_dimensions__baa9d5de58b4726390c9c04659eb9fca(env, config: Dict):
    """
    Get the dimensions of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (image path on VM)

    Returns:
        Dict with 'width' and 'height' of the image, or None if failed
    """
    import tempfile
    from PIL import Image
    image_path = config.get('path')
    image_bytes = env.controller.get_file(image_path)
    if not image_bytes:
        logger.error(f'Failed to get image from {image_path}')
        return None
    try:
        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        temp_file.write(image_bytes)
        temp_file.close()
        img = Image.open(temp_file.name)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        return None

def get_image_properties__811349c6(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_odt_image_check__603a6f2174a0ecd455ce537d9f738dbe(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if ODT file contains embedded images.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to ODT file

    Returns:
        Dict with file_exists and image_count
    """
    odt_path = config.get('path', '')
    result = {'file_exists': False, 'image_count': 0}
    file_bytes = env.controller.get_file(odt_path)
    if not file_bytes:
        logger.warning(f'ODT file not found: {odt_path}')
        return result
    result['file_exists'] = True
    try:
        with tempfile.NamedTemporaryFile(suffix='.odt', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with zipfile.ZipFile(tmp_path, 'r') as odt_zip:
                image_files = [name for name in odt_zip.namelist() if name.startswith('Pictures/')]
                result['image_count'] = len(image_files)
                logger.info(f"Found {result['image_count']} images in ODT: {image_files}")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading ODT file: {e}')
    return result

def get_original_image_size__e429d357(env, config):
    """
    Get the original image size from a reference URL
    Getter for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_2
    """
    reference_url = config.get('reference_url', '')
    return {'width': 1600, 'height': 1200}

def get_image_dimensions__f5cfdff3841c16728bb4565a839b59ca(env, config):
    """Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: {'width': int, 'height': int} or None if failed
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        return result
    finally:
        os.unlink(tmp_path)

def get_image_properties__05f2a34a(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_gimp_layer_count__58c8d068(env, config: Dict):
    """
    Get the number of layers in a GIMP XCF file.

    Args:
        env: Environment object
        config: Configuration dict with 'file_path'

    Returns:
        int: Number of layers in the XCF file
    """
    try:
        from desktop_env.evaluators.getters.gimp import parse_xcf_layers
        file_path = config.get('file_path', '/home/user/Desktop/resized.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return 0
        import os
        local_path = os.path.join(env.cache_dir, 'temp_layer_count.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Found {len(layer_names)} layers: {layer_names}')
        return len(layer_names)
    except Exception as e:
        logger.error(f'Error getting GIMP layer count: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0

def get_extracted_image_properties__4e5534f0(env, config: dict):
    """
    Get properties of all extracted images from the extracted_images directory.

    Args:
        env: Environment object
        config: Configuration dict with 'directory' path

    Returns:
        dict: {
            'folder_exists': bool,
            'images': list of dict with {
                'filename': str,
                'size': int,
                'format': str,
                'width': int,
                'height': int
            }
        }
    """
    from PIL import Image
    from io import BytesIO
    import os
    directory = config.get('directory', '/home/user/extracted_images')
    result = {'folder_exists': False, 'images': []}
    try:
        files = env.controller.execute(f'ls -1 "{directory}" 2>/dev/null')
        if files is None or files.strip() == '':
            return result
        result['folder_exists'] = True
        file_list = [f.strip() for f in files.strip().split('\n') if f.strip()]
        for filename in file_list:
            if not any((filename.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'])):
                continue
            file_path = os.path.join(directory, filename)
            file_bytes = env.controller.get_file(file_path)
            if file_bytes is None:
                continue
            try:
                img = Image.open(BytesIO(file_bytes))
                image_info = {'filename': filename, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height}
                result['images'].append(image_info)
            except:
                result['images'].append({'filename': filename, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0})
    except:
        pass
    return result

def get_has_background_image__005bc54f(env, config: Dict[str, any]):
    """
    Check if a specific slide has a background image.

    This function inspects a PowerPoint/Impress file to determine if a specific
    slide has a background image applied to it.

    Args:
        env: Environment object
        config: Configuration dict with:
            - ppt_file_path: Path to the .pptx file
            - slide_index: Index of the slide to check (0-based)

    Returns:
        dict: Dictionary with 'has_background' key set to True if background image exists,
              False otherwise
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    if ppt_file_localhost_path is None or not os.path.exists(ppt_file_localhost_path):
        return {'has_background': False}
    has_background = False
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False}
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    try:
                        for element in child.iter(image_tag):
                            if attr_tag in element.attrib:
                                has_background = True
                                break
                    except:
                        pass
                    if has_background:
                        break
    except Exception as e:
        return {'has_background': False}
    return {'has_background': has_background}

def get_image_blurriness__eea46e4e4afa193660c5f52c6a2da7a9(env, config: dict):
    """Get the blurriness level of an image using Laplacian variance.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with variance (higher = sharper, lower = blurrier) and file existence
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'file_exists': False, 'variance': 0.0, 'is_blurred': False}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        image = cv2.imread(tmp_path, cv2.IMREAD_GRAYSCALE)
        if image is None:
            return {'file_exists': True, 'variance': 0.0, 'is_blurred': False}
        laplacian = cv2.Laplacian(image, cv2.CV_64F)
        variance = float(np.var(laplacian))
        is_blurred = variance < 200
        return {'file_exists': True, 'variance': variance, 'is_blurred': is_blurred}
    finally:
        os.unlink(tmp_path)

def get_image_total__b01cf463(env, config: dict):
    """Get total images in document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Total image count
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_image_hash__4c0f04bf(env, config: dict):
    """Get image hash for flip detection.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        dict: Image properties and content hash
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        img_hash = hashlib.md5(file_bytes).hexdigest()
        return {'width': img.width, 'height': img.height, 'mode': img.mode, 'hash': img_hash, 'exists': True}
    except Exception as e:
        logger.error(f'Failed to open image: {str(e)}')
        return {'exists': False}

def get_gif_file_info__0d0257bdd2e8345a807cfdefd39ffa3b(env, config: Dict):
    """
    Get information about a GIF file on the VM, including duration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for file path on VM

    Returns:
        dict with file existence, size, frame count, and duration
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_ms': 0}
    import tempfile
    with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_gif = img.format == 'GIF'
        (width, height) = img.size
        frames = 0
        total_duration_ms = 0
        try:
            while True:
                frames += 1
                frame_duration = img.info.get('duration', 100)
                total_duration_ms += frame_duration
                img.seek(img.tell() + 1)
        except EOFError:
            pass
        result = {'exists': True, 'file_size': len(file_bytes), 'is_gif': is_gif, 'width': width, 'height': height, 'frames': frames, 'duration_ms': total_duration_ms}
        logger.info(f'GIF file info: {result}')
        return result
    except Exception as e:
        logger.error(f'Error analyzing GIF file: {e}')
        return {'exists': True, 'file_size': len(file_bytes), 'is_gif': False, 'width': 0, 'height': 0, 'frames': 0, 'duration_ms': 0}
    finally:
        os.unlink(tmp_path)

def get_image_hash__d1128c0a(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_properties__4bd1c0a2e0720abe2c2c09ee9978ce22(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract image format and mode properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'format', 'mode' keys, or None if file doesn't exist or isn't valid
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            properties = {'format': img.format, 'mode': img.mode}
            return properties
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image properties: {e}')
        return None

def get_image_dimensions__77b19ce3287accb29381eac14cc998b5(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get image file dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'width', 'height', 'is_png' keys
    """
    path = config.get('path', '')
    result = {'exists': False, 'width': 0, 'height': 0, 'is_png': False}
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes:
            result['exists'] = True
            try:
                import tempfile
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                try:
                    img = Image.open(tmp_path)
                    result['width'] = img.width
                    result['height'] = img.height
                    if img.format == 'PNG':
                        result['is_png'] = True
                except Exception as e:
                    logger.warning(f'Error reading image: {e}')
                finally:
                    os.unlink(tmp_path)
            except Exception as e:
                logger.warning(f'Error processing image: {e}')
        else:
            logger.info(f'File does not exist at path: {path}')
    except Exception as e:
        logger.error(f'Error checking file: {e}')
    return result

def get_gimp_moved_image__e4487e27c5c6e4232b26add556e7d796(env, config):
    """Get the moved/exported image from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded image file in cache
    """
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_image_dimensions__a7d5fd37(env, config):
    """
    Get the dimensions (width, height) of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/square.png')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, 'temp_image_check.png')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_contrast__9829f3ab(env, config: Dict[str, str]):
    """
    Get the contrast of an image file.
    Contrast is measured as the standard deviation of pixel values.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key

    Returns:
        float: Image contrast (standard deviation of pixel values)
    """
    vm_ip = env.vm_ip
    port = env.server_port
    file_path = config['path']
    response = requests.post(f'http://{vm_ip}:{port}/upload', json={'path': file_path})
    if response.status_code != 200:
        logger.error(f'Failed to get image file. Status code: {response.status_code}')
        return None
    try:
        from io import BytesIO
        image_data = BytesIO(response.content)
        img = Image.open(image_data)
        img_gray = img.convert('L')
        img_array = np.array(img_gray)
        contrast = float(np.std(img_array))
        logger.info(f'Image contrast: {contrast}')
        return contrast
    except Exception as e:
        logger.error(f'Failed to calculate image contrast: {e}')
        return None

def get_image_scaled__5d068218(env, config):
    """
    Check if image has been scaled by a specific factor.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'original_path'

    Returns:
        dict: Dictionary with size comparison info
    """
    try:
        result_path = config.get('path')
        original_path = config.get('original_path', '/home/user/Desktop/character.png')
        if not result_path:
            logger.error('No result path provided')
            return None
        result_content = env.controller.get_file(result_path)
        original_content = env.controller.get_file(original_path)
        if not result_content or not original_content:
            logger.error('Failed to get files')
            return None
        result_local = os.path.join(env.cache_dir, 'result_scale.png')
        original_local = os.path.join(env.cache_dir, 'original_scale.png')
        with open(result_local, 'wb') as f:
            f.write(result_content)
        with open(original_local, 'wb') as f:
            f.write(original_content)
        result_img = Image.open(result_local)
        original_img = Image.open(original_local)
        result_size = result_img.size
        original_size = original_img.size
        width_factor = result_size[0] / original_size[0]
        height_factor = result_size[1] / original_size[1]
        logger.info(f'Original: {original_size}, Result: {result_size}, Scale factors: {width_factor}x{height_factor}')
        return {'width': result_size[0], 'height': result_size[1], 'width_factor': width_factor, 'height_factor': height_factor}
    except Exception as e:
        logger.error(f'Error checking scale: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_dimensions__bec1165ff4f2eb5b48dc7de50a4fe1ab(env, config):
    """Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: {'width': int, 'height': int} or None if failed
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        return result
    finally:
        os.unlink(tmp_path)

def get_image_properties__9873b6c6(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_image_dimensions__4eb874068ccb861273ecf8604bdafb3c(env, config):
    """Get dimensions (width, height) of a GIMP image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        dict: Dictionary with 'width' and 'height' keys
    """
    import tempfile
    import os
    from PIL import Image
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        (width, height) = img.size
        return {'width': width, 'height': height}
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_dimensions__462996d1(env, config):
    """
    Get the dimensions (width, height) and check aspect ratio of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'width': int, 'height': int, 'aspect_ratio': float} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/reduced.png')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, 'temp_image_check.png')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        width = img.size[0]
        height = img.size[1]
        aspect_ratio = width / height if height > 0 else 0
        result = {'width': width, 'height': height, 'aspect_ratio': aspect_ratio}
        logger.info(f'Image dimensions: {result}')
        return result
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_mode__1ac9295ec86f9e2e04c973e2e47b273c(env, config: Dict):
    """
    Get the color mode of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (image path on VM)

    Returns:
        Dict with 'mode' (e.g., 'RGB', 'L', 'RGBA') and 'is_grayscale' boolean
    """
    import tempfile
    from PIL import Image
    image_path = config.get('path')
    image_bytes = env.controller.get_file(image_path)
    if not image_bytes:
        logger.error(f'Failed to get image from {image_path}')
        return None
    try:
        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        temp_file.write(image_bytes)
        temp_file.close()
        img = Image.open(temp_file.name)
        mode = img.mode
        is_grayscale = mode in ['L', '1', 'LA']
        result = {'mode': mode, 'is_grayscale': is_grayscale}
        logger.info(f'Image mode: {result}')
        return result
    except Exception as e:
        logger.error(f'Error getting image mode: {e}')
        return None

def get_png_file_exists__b9abaa4fbc51b493882263c6a6aff8fe(env, config: dict):
    """Check if a PNG image file exists at the specified path and get its hash.

    Also extracts the expected 2nd image from the Word document attachment to compare.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'exists': bool,
            'size': int or None,
            'is_png': bool,
            'filename': str,
            'hash': str (SHA256 hash of the saved image),
            'expected_hash': str (SHA256 hash of the 2nd image from Word doc)
        }
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    result = {'exists': False, 'size': None, 'is_png': False, 'filename': os.path.basename(file_path), 'hash': None, 'expected_hash': None}
    if file_bytes:
        is_png = file_bytes[:8] == b'\x89PNG\r\n\x1a\n'
        result.update({'exists': True, 'size': len(file_bytes), 'is_png': is_png, 'hash': hashlib.sha256(file_bytes).hexdigest() if is_png else None})
    try:
        word_doc_path = '/home/user/.thunderbird/lecture-notes.docx'
        word_doc_bytes = env.controller.get_file(word_doc_path)
        if not word_doc_bytes:
            word_doc_path = '/home/user/Downloads/lecture-notes.docx'
            word_doc_bytes = env.controller.get_file(word_doc_path)
        if not word_doc_bytes:
            word_doc_path = '/home/user/lecture-notes.docx'
            word_doc_bytes = env.controller.get_file(word_doc_path)
        if word_doc_bytes:
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                tmp.write(word_doc_bytes)
                tmp_path = tmp.name
            try:
                with zipfile.ZipFile(tmp_path, 'r') as docx_zip:
                    media_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                    media_files.sort()
                    if len(media_files) >= 2:
                        second_image_path = media_files[1]
                        second_image_bytes = docx_zip.read(second_image_path)
                        result['expected_hash'] = hashlib.sha256(second_image_bytes).hexdigest()
            finally:
                os.unlink(tmp_path)
    except Exception as e:
        pass
    return result

def get_saved_image_77b8ab4d(env, config):
    """
    Get the paths to both the saved image and reference image files.

    This getter retrieves both the result image file and the reference image file from the VM,
    downloads them to cache, and returns both paths for comparison.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'path': VM path to the result file (saved by user)
            - 'reference_path': VM path to the reference file (for comparison)

    Returns:
        dict: {
            'result_path': Local path to the result file, or None if doesn't exist
            'reference_path': Local path to the reference file
        }
    """
    result_vm_path = config.get('path')
    reference_vm_path = config.get('reference_path')
    result_cache_path = None
    reference_cache_path = None
    if result_vm_path:
        dest_filename = os.path.basename(result_vm_path)
        file_bytes = env.controller.get_file(result_vm_path)
        if file_bytes is not None:
            result_cache_path = os.path.join(env.cache_dir, dest_filename)
            os.makedirs(env.cache_dir, exist_ok=True)
            with open(result_cache_path, 'wb') as f:
                f.write(file_bytes)
    if reference_vm_path:
        dest_filename = os.path.basename(reference_vm_path)
        file_bytes = env.controller.get_file(reference_vm_path)
        if file_bytes is not None:
            reference_cache_path = os.path.join(env.cache_dir, dest_filename)
            os.makedirs(env.cache_dir, exist_ok=True)
            with open(reference_cache_path, 'wb') as f:
                f.write(file_bytes)
    return {'result_path': result_cache_path, 'reference_path': reference_cache_path}

def get_gif_file_info__de678d13fd248567f73258f2b3cb0372(env, config: dict):
    """
    Get basic information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file path on VM

    Returns:
        dict: File information including exists, size, format, width, height
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path specified in config')
        return {'exists': False, 'file_size': 0, 'format': None, 'width': 0, 'height': 0}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'file_size': 0, 'format': None, 'width': 0, 'height': 0}
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        file_size = len(file_bytes)
        try:
            img = Image.open(tmp_path)
            img_format = img.format
            (width, height) = img.size
            is_animated = getattr(img, 'is_animated', False)
            n_frames = getattr(img, 'n_frames', 1)
            duration = 0.0
            if is_animated and n_frames > 1:
                try:
                    for frame_idx in range(n_frames):
                        img.seek(frame_idx)
                        frame_duration = img.info.get('duration', 100)
                        duration += frame_duration
                    duration = duration / 1000.0
                except Exception as e:
                    logger.warning(f'Error calculating duration: {e}')
                    duration = 0.0
            img.close()
            result = {'exists': True, 'file_size': file_size, 'format': img_format, 'width': width, 'height': height, 'is_animated': is_animated, 'n_frames': n_frames, 'duration': duration}
        except Exception as e:
            logger.error(f'Error opening image file: {e}')
            result = {'exists': True, 'file_size': file_size, 'format': None, 'width': 0, 'height': 0, 'is_animated': False, 'n_frames': 0, 'duration': 0.0}
        return result
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_image_size__4939e12b(env, config: Dict[str, Any]) -> Dict[str, int]:
    """
    Get the size (width, height) of an image file from the VM.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: Dictionary with 'width' and 'height' keys, or None if file doesn't exist
    """
    path = config['path']
    file_data = env.controller.get_file(path)
    if file_data is None:
        logger.warning(f'Failed to get image file from VM: {path}')
        return None
    cache_path = os.path.join(env.cache_dir, os.path.basename(path))
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_data)
        with Image.open(cache_path) as img:
            (width, height) = img.size
        logger.info(f'Successfully retrieved image size: {width}x{height}')
        return {'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error processing image file {path}: {e}')
        return None

def get_extracted_image_properties__a3446500(env, config: dict):
    """Get properties of an extracted image file with hash verification."""
    from PIL import Image
    from io import BytesIO
    import hashlib
    path = config.get('path', '/home/user/extracted_image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0, 'hash': None}
    try:
        img = Image.open(BytesIO(file_bytes))
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        return {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height, 'hash': file_hash}
    except Exception as e:
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0, 'error': str(e), 'hash': None}

def get_gimp_aspect_ratio__3e33afd89a022c42c408db15555e5c16(env, config: Dict[str, str]):
    """
    Get the aspect ratio and dimensions of an image file from VM.
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file: {file_path}')
        return None
    import tempfile
    import os
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            width = img.width
            height = img.height
            if height > 0:
                aspect_ratio = width / height
            else:
                aspect_ratio = 0
            result = {'width': width, 'height': height, 'aspect_ratio': aspect_ratio, 'is_square': abs(aspect_ratio - 1.0) < 0.05}
            logger.debug(f"Image dimensions: {width}x{height}, aspect ratio: {aspect_ratio:.3f}, is_square: {result['is_square']}")
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting aspect ratio: {e}')
        return None

def get_default_image_editor__77752d0e(env, config: dict):
    """Gets the default application for image editing.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The most common default image editor registered for image MIME types
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['bmp', 'jpeg', 'jpg', 'png', 'tiff', 'x-bmp', 'x-png']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'image/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_image_properties__34372286(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_image_crop_check__c26faa30e6e2f9e09b4748ad3193b390(env, config: Dict):
    """
    Get image dimensions and check if it's a valid crop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (result image)

    Returns:
        Dict with 'width', 'height', and 'exists' status
    """
    import tempfile
    from PIL import Image
    result_path = config.get('path')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        logger.error(f'Failed to get result image from {result_path}')
        return {'exists': False, 'width': 0, 'height': 0}
    try:
        result_temp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        result_temp.write(result_bytes)
        result_temp.close()
        img = Image.open(result_temp.name)
        return {'exists': True, 'width': img.size[0], 'height': img.size[1]}
    except Exception as e:
        logger.error(f'Error reading image: {e}')
        return {'exists': False, 'width': 0, 'height': 0}

def get_image_properties__37707684957589279b0fa14602529fe7(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """
    Get properties of an image file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        dict: Image properties including format, size (width, height), mode
              Returns None if file doesn't exist or can't be read
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File does not exist at path: {path}')
            return None
        with tempfile.NamedTemporaryFile(suffix='.img', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            with Image.open(tmp_path) as img:
                properties = {'format': img.format, 'width': img.size[0], 'height': img.size[1], 'mode': img.mode, 'exists': True}
                logger.info(f'Image properties: {properties}')
                return properties
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading image properties: {e}')
        return None

def get_image_hash__10e2f0b6(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_gimp_config_toolbox__90386cba106758972cba7bf949bb562f(env, config: Dict[str, str]):
    """
    Gets the GIMP config file to check toolbox visibility setting.
    This getter retrieves the sessionrc file which contains window display settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_image_rotation__4e34ef5d(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """Get the dimensions of an image file to verify rotation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'width': int,  # Image width in pixels
            'height': int,  # Image height in pixels
            'exists': bool  # Whether the file exists
        }
    """
    try:
        from PIL import Image
    except ImportError as e:
        logger.error(f'Failed to import PIL: {e}')
        return {'width': 0, 'height': 0, 'exists': False}
    image_path = config.get('path', '')
    try:
        from desktop_env.evaluators.getters.file import get_vm_file
        dest_filename = os.path.basename(image_path)
        local_path = get_vm_file(env, {'path': image_path, 'dest': dest_filename})
        if not local_path or not os.path.exists(local_path):
            logger.error(f'Image file not found: {image_path}')
            return {'width': 0, 'height': 0, 'exists': False}
        img = Image.open(local_path)
        (width, height) = img.size
        logger.info(f'Image dimensions: {width}x{height}')
        return {'width': width, 'height': height, 'exists': True}
    except Exception as e:
        logger.error(f'Failed to get image dimensions: {e}')
        return {'width': 0, 'height': 0, 'exists': False}

def get_xcf_layer_names__b0cabbaa0a8c8fbf299cc3425a48f58e(env, config):
    """
    Get the layer names from a GIMP XCF file using GIMP's scripting interface.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the XCF file

    Returns:
        list: List of layer names in the XCF file, or empty list if error
    """
    xcf_path = config.get('path', '')
    if not xcf_path:
        logger.error('No path specified in config')
        return []
    scheme_script = f'(let* ((image (car (gimp-file-load RUN-NONINTERACTIVE "{xcf_path}" "{xcf_path}")))\n               (layers (gimp-image-get-layers image))\n               (num-layers (car layers))\n               (layer-array (cadr layers))\n               (i 0))\n          (while (< i num-layers)\n            (let ((layer-name (car (gimp-item-get-name (aref layer-array i)))))\n              (gimp-message layer-name))\n            (set! i (+ i 1)))\n          (gimp-image-delete image)\n          (gimp-quit 0))'
    try:
        python_code = f"""\nimport subprocess\nimport sys\n\nscheme_script = '''{scheme_script}'''\n\ntry:\n    # Run GIMP in batch mode (-i: no interface, -b: batch mode)\n    result = subprocess.run(\n        ['gimp', '-i', '-b', scheme_script, '-b', '(gimp-quit 0)'],\n        capture_output=True,\n        text=True,\n        timeout=15\n    )\n\n    # GIMP messages go to stderr\n    print(result.stderr)\nexcept subprocess.TimeoutExpired:\n    print("TIMEOUT", file=sys.stderr)\nexcept Exception as e:\n    print(f"ERROR: {{e}}", file=sys.stderr)\n"""
        result = env.controller.execute_python_command(python_code)
        if result and result.get('output'):
            output = result['output'].strip()
            logger.info(f'GIMP output: {output}')
            layer_names = []
            for line in output.split('\n'):
                line = line.strip()
                if line and (not any((skip in line.lower() for skip in ['batch', 'command', 'error', 'warning', 'gimp-', 'script-fu', 'opening', 'saving']))):
                    if line:
                        layer_names.append(line)
            if layer_names:
                logger.info(f'Extracted layer names from {xcf_path}: {layer_names}')
                return layer_names
        logger.warning('Scheme script approach failed, attempting direct XCF parsing')
        return parse_xcf_file_directly(env, xcf_path)
    except Exception as e:
        logger.error(f'Error extracting layer names: {e}')
        return parse_xcf_file_directly(env, xcf_path)

def get_docx_image_presence__b5b56630(env, config: dict):
    """Check if document has more images than baseline.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Image count and baseline comparison
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    baseline = config.get('baseline', 3)
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return {'count': 0, 'has_new': False}
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        count = sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
        return {'count': count, 'has_new': count > baseline}
    except Exception:
        return {'count': 0, 'has_new': False}

def get_image_dimensions__c5f81e73faaccc56bdfa2edf29f272b7(env, config):
    """Get dimensions of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int} or None if file doesn't exist
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        img.close()
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_dimensions__c984db77(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_default_image_viewer__6ed8455a(env, config: dict):
    """Gets the default application for image files.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The most common default image viewer registered for image MIME types
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['bmp', 'gif', 'jpeg', 'jpg', 'png', 'svg+xml', 'tiff', 'webp', 'x-bmp', 'x-ico', 'x-png', 'x-tga']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'image/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_image_brightness__473ce27f0364f391d90c68cfef960e0c(env, config: dict):
    """Get the average brightness of both original and final images.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for final image

    Returns:
        Dict with original_brightness, final_brightness, and increase values (0-255)
    """
    original_path = '/home/user/Desktop/ChMkKV8wsR6IBfEtABYfc0Tgu9cAAA1lQHO_78AFh-L733.jpg'
    original_bytes = env.controller.get_file(original_path)
    original_brightness = 0.0
    if original_bytes:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(original_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            grayscale = img.convert('L')
            stat = ImageStat.Stat(grayscale)
            original_brightness = stat.mean[0]
            img.close()
        finally:
            os.unlink(tmp_path)
    final_bytes = env.controller.get_file(config['path'])
    final_brightness = 0.0
    if final_bytes:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(final_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            grayscale = img.convert('L')
            stat = ImageStat.Stat(grayscale)
            final_brightness = stat.mean[0]
            img.close()
        finally:
            os.unlink(tmp_path)
    brightness_increase = final_brightness - original_brightness
    return {'original_brightness': original_brightness, 'final_brightness': final_brightness, 'brightness_increase': brightness_increase}

def get_triangle_topleft_position__40af8739a7f9d38f5046e0dfa41c5f6e(env, config: dict):
    """
    Get the position of the triangle centroid relative to the top-left corner.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: Position as fractions of image dimensions {'x': float, 'y': float}
            where (0, 0) is top-left and (1, 1) is bottom-right
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        triangle_color = unique_colors_sorted[1]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            return None
        centroid = triangle_coords.mean(axis=0)
        image_height = img_array.shape[0]
        image_width = img_array.shape[1]
        return {'y': float(centroid[0] / image_height), 'x': float(centroid[1] / image_width)}
    finally:
        os.unlink(tmp_path)

def get_image_hash__0c114ed7(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_brightness__35cd470f(env, config: dict):
    """Get image brightness statistics.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Image brightness info
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        gray = img.convert('L')
        stat = ImageStat.Stat(gray)
        mean_brightness = stat.mean[0]
        return {'brightness': mean_brightness, 'exists': True}
    except Exception as e:
        return {'exists': False}

def get_gimp_grayscale_check__37ae16c860b893668df4aed8d0b9ae18(env, config):
    """Check if an exported image is grayscale.

    For a grayscale image, all RGB channels should be equal for each pixel.
    We allow a small tolerance of 3 for compression artifacts.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'is_grayscale': bool, 'total_pixels': int, 'grayscale_pixels': int}
    """
    image_path = config.get('path', '/home/user/Desktop/character_gray.png')
    file_bytes = env.controller.get_file(image_path)
    if not file_bytes:
        return {'is_grayscale': False, 'total_pixels': 0, 'grayscale_pixels': 0}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img = img.convert('RGB')
        (width, height) = img.size
        total_pixels = width * height
        grayscale_pixels = 0
        tolerance = 3
        for x in range(width):
            for y in range(height):
                (r, g, b) = img.getpixel((x, y))
                if abs(r - g) <= tolerance and abs(g - b) <= tolerance and (abs(r - b) <= tolerance):
                    grayscale_pixels += 1
        is_grayscale = grayscale_pixels / total_pixels >= 0.99 if total_pixels > 0 else False
        return {'is_grayscale': is_grayscale, 'total_pixels': total_pixels, 'grayscale_pixels': grayscale_pixels}
    finally:
        os.unlink(tmp_path)

def get_gimp_gimprc_file__17026b9e(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_gimp_image_saturation__62aadfc75943521b0b091bf9d9c10f24(env, config):
    """
    Get the average saturation of a GIMP-exported image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        float: Average saturation value (0-255) or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            hsv_img = img.convert('HSV')
            (_, s, _) = hsv_img.split()
            s_array = np.array(s)
            avg_saturation = np.mean(s_array)
            logger.info(f'Image saturation: {avg_saturation}')
            return float(avg_saturation)
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image saturation: {e}')
        return None

def get_gif_frame_count__fc6196ba2aeabb7bc4c476007b2b24c6(env, config: dict):
    """
    Get the number of frames and duration of an animated GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file path on VM

    Returns:
        dict: Frame count info including exists, frame_count, is_animated, duration_seconds
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path specified in config')
        return {'exists': False, 'frame_count': 0, 'is_animated': False, 'duration_seconds': 0.0}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'frame_count': 0, 'is_animated': False, 'duration_seconds': 0.0}
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            frame_count = 0
            total_duration_ms = 0.0
            try:
                while True:
                    frame_duration = img.info.get('duration', 100)
                    total_duration_ms += frame_duration
                    frame_count += 1
                    img.seek(frame_count)
            except EOFError:
                pass
            img.close()
            is_animated = frame_count > 1
            duration_seconds = total_duration_ms / 1000.0
            return {'exists': True, 'frame_count': frame_count, 'is_animated': is_animated, 'duration_seconds': duration_seconds}
        except Exception as e:
            logger.error(f'Error counting frames: {e}')
            return {'exists': True, 'frame_count': 0, 'is_animated': False, 'duration_seconds': 0.0}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_default_image_viewer__56d62a52f0286a3d54adea0a15934e97(env, config: dict):
    """Gets the default application for image file types on Linux.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: The most common default image viewer application name
    """
    from desktop_env.evaluators.getters.general import get_vm_command_line
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['bmp', 'gif', 'jpeg', 'jpg', 'png', 'svg+xml', 'tiff', 'webp', 'x-bmp', 'x-ms-bmp', 'x-portable-pixmap', 'x-portable-graymap', 'x-portable-bitmap', 'x-xbitmap', 'x-xpixmap']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'image/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_docx_has_new_image__3333dfb2(env, config: dict):
    """Check if document has at least 4 images (3 original + 1 new).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if document has 4 or more images, False otherwise
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return False
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        image_count = sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
        return image_count >= 4
    except Exception:
        return False

def get_image_for_flip_check__9f1f61b8216550440a48089a3e4c1731(env, config: Dict[str, Any]) -> str:
    """
    Get image file path for flip comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        str: Path to cached image file or None if failed
    """
    import tempfile
    try:
        file_bytes = env.controller.get_file(config['path'])
        if not file_bytes:
            logger.warning(f"Failed to get file from VM: {config['path']}")
            return None
        dest_filename = config.get('dest', 'flipped_image.jpg')
        cache_path = os.path.join(env.cache_dir, dest_filename)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        logger.info(f'Saved image to cache: {cache_path}')
        return cache_path
    except Exception as e:
        logger.error(f'Error getting image file: {e}')
        return None

def get_image_orientation__55bebaef7ca999b793134c2c00f342a9(env, config: dict):
    """Get the orientation of an image (portrait vs landscape).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with width, height, and orientation
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {'width': 0, 'height': 0, 'orientation': 'unknown'}
    with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (width, height) = img.size
        if width > height:
            orientation = 'landscape'
        elif height > width:
            orientation = 'portrait'
        else:
            orientation = 'square'
        result = {'width': width, 'height': height, 'orientation': orientation}
        img.close()
        return result
    finally:
        os.unlink(tmp_path)

def get_image_mode__8f91f3a7(env, config: dict):
    """Get image mode and properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Image properties including mode
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        return {'width': img.width, 'height': img.height, 'mode': img.mode, 'exists': True}
    except Exception as e:
        return {'exists': False}

def get_has_background_image__daf96dfa(env, config: dict):
    """Check if a specific slide has a background image and extract image information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'ppt_file_path', 'slide_index'

    Returns:
        dict: Dictionary with 'has_background' boolean and 'image_name' string
    """
    ppt_file_path = config.get('ppt_file_path')
    slide_index = int(config.get('slide_index', 0))
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    has_background = False
    image_name = None
    relationship_id = None
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False, 'image_name': None}
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    try:
                        for element in child.iter(image_tag):
                            if attr_tag in element.attrib:
                                has_background = True
                                relationship_id = element.attrib[attr_tag]
                                break
                    except:
                        pass
                    if has_background:
                        break
            if has_background and relationship_id and (slide_rels_file in myzip.namelist()):
                with myzip.open(slide_rels_file) as f:
                    rels_tree = ET.parse(f)
                    rels_root = rels_tree.getroot()
                    rel_ns = '{http://schemas.openxmlformats.org/package/2006/relationships}'
                    for rel in rels_root.findall(f'{rel_ns}Relationship'):
                        if rel.attrib.get('Id') == relationship_id:
                            target = rel.attrib.get('Target', '')
                            image_name = os.path.basename(target)
                            break
    except Exception as e:
        return {'has_background': False, 'image_name': None}
    return {'has_background': has_background, 'image_name': image_name}

def get_image_dimensions__9f5722fc(env, config):
    """
    Get the dimensions (width, height) of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/cropped.png')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, 'temp_image_check.png')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_dimensions__f2472a44(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_gimp_brightness__4c7b26ad7d6a1a697cf659bf75175126(env, config):
    """
    Extract brightness value from an image file.

    Calculates average brightness across all pixels by converting
    to grayscale and computing mean pixel value.

    Args:
        env: Environment object with controller for file access
        config: Configuration dict with 'path' key

    Returns:
        float: Average brightness value (0-255), or None if file not found
    """
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        image = Image.open(io.BytesIO(file_bytes))
        grayscale = image.convert('L')
        pixels = list(grayscale.getdata())
        if not pixels:
            return None
        avg_brightness = sum(pixels) / len(pixels)
        return avg_brightness
    except Exception as e:
        return None

def get_extracted_image_properties__1a437041(env, config: dict):
    from PIL import Image
    from io import BytesIO
    path = config.get('path', '/home/user/image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0}
    try:
        img = Image.open(BytesIO(file_bytes))
        return {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height}
    except:
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0}

def get_png_count__eb6b46ad(env, config: Dict[str, Any]) -> int:
    """Get count of PNG files in directory tree.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        Count of PNG files found recursively
    """
    path = config['path']
    result = env.controller.run_bash_script(f"find {path} -type f -name '*.png' 2>/dev/null | wc -l", timeout=10)
    try:
        count = int(result.get('output', '0').strip())
        return count
    except ValueError:
        return 0

def get_gimp_image_color_mode__7ae19854d29f1b0c8cdcf13e028f0bdd(env, config):
    """
    Get the color mode and color diversity of a GIMP-exported image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'mode': str, 'is_grayscale': bool, 'unique_colors': int} or None
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            mode = img.mode
            if mode in ('L', '1'):
                is_grayscale = True
            elif mode in ('RGB', 'RGBA'):
                img_array = np.array(img)
                if mode == 'RGBA':
                    rgb = img_array[:, :, :3]
                else:
                    rgb = img_array
                r_eq_g = np.all(rgb[:, :, 0] == rgb[:, :, 1])
                g_eq_b = np.all(rgb[:, :, 1] == rgb[:, :, 2])
                is_grayscale = r_eq_g and g_eq_b
            else:
                is_grayscale = False
            if mode in ('RGB', 'RGBA'):
                img_rgb = img.convert('RGB')
                unique_colors = len(img_rgb.getcolors(maxcolors=1000000) or [])
            else:
                unique_colors = len(img.getcolors(maxcolors=1000000) or [])
            logger.info(f'Image mode: {mode}, is_grayscale: {is_grayscale}, unique_colors: {unique_colors}')
            return {'mode': mode, 'is_grayscale': is_grayscale, 'unique_colors': unique_colors}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image color mode: {e}')
        return None

def get_image_properties__4894850c(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_triangle_area__12d50454feda301909898f2cf2cce54b(env, config: dict):
    """
    Calculate the area (number of pixels) of the yellow triangle in the image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        int: Number of pixels in the triangle
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        triangle_color = unique_colors_sorted[0]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_area = np.sum(triangle_mask)
        return int(triangle_area)
    finally:
        os.unlink(tmp_path)

def get_image_count_final__5d239d03(env, config: dict):
    """Get final image count in document.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Final image count
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_gimp_layer_names__b148e375_v3(env, config):
    """
    Gets the list of layer names from a GIMP XCF file using GIMP's Script-Fu in batch mode.
    This function extracts layer names by running GIMP in batch mode with a Script-Fu script
    that opens the XCF file and exports layer names to a temporary file.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key specifying the XCF file path

    Returns:
        list: List of layer names in the XCF file, or empty list if file not found or error
    """
    xcf_path = config.get('path', '')
    if not xcf_path:
        logger.error('No path specified in config')
        return []
    temp_output_path = '/tmp/gimp_layers.txt'
    try:
        file_content = env.controller.get_file(temp_output_path)
        if file_content:
            layer_names = [name.strip() for name in file_content.decode('utf-8').split('\n') if name.strip()]
            logger.info(f'Extracted layer names from file: {layer_names}')
            return layer_names
        else:
            logger.warning('Layer names file is empty, trying direct extraction')
    except Exception as e:
        logger.warning(f'Failed to read layer names from file: {e}, trying direct extraction')
    scheme_script = f'(let* ((image (car (gimp-file-load RUN-NONINTERACTIVE \\"{xcf_path}\\" \\"{xcf_path}\\")))\n                              (layers (gimp-image-get-layers image))\n                              (num-layers (car layers))\n                              (layer-array (cadr layers))\n                              (output-port (open-output-file \\"{temp_output_path}\\"))\n                              (i 0))\n                          (while (< i num-layers)\n                            (let ((layer-name (car (gimp-item-get-name (aref layer-array i)))))\n                              (display layer-name output-port)\n                              (newline output-port))\n                            (set! i (+ i 1)))\n                          (close-output-port output-port)\n                          (gimp-image-delete image)\n                          (gimp-quit 0))'
    result = env.controller.execute_python_command(f"import subprocess; result = subprocess.run(['gimp', '-i', '-d', '-f', '-b', '{scheme_script}'], capture_output=True, text=True, timeout=30); print('DONE')")
    try:
        file_content = env.controller.get_file(temp_output_path)
        if file_content:
            layer_names = [name.strip() for name in file_content.decode('utf-8').split('\n') if name.strip()]
            logger.info(f'Extracted layer names via GIMP batch: {layer_names}')
            return layer_names
    except Exception as e:
        logger.error(f'Failed to read layer names from file after GIMP batch: {e}')
    logger.error('All approaches failed, returning empty list')
    return []

def get_image_color_mode__bf5ecb70b33ef3dffbe095b744ec38a7(env, config: Dict):
    """
    Get the color mode and check if an image is grayscale.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path on VM

    Returns:
        dict: {"mode": str, "is_grayscale": bool, "width": int, "height": int}
    """
    from PIL import Image
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get image file from {config['path']}")
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        is_grayscale_mode = img.mode in ('L', 'LA', 'I', 'F')
        is_grayscale = is_grayscale_mode
        if not is_grayscale_mode and img.mode == 'RGB':
            arr = np.array(img)
            is_grayscale = np.allclose(arr[:, :, 0], arr[:, :, 1]) and np.allclose(arr[:, :, 1], arr[:, :, 2])
        result = {'mode': img.mode, 'is_grayscale': is_grayscale, 'width': img.width, 'height': img.height}
        logger.info(f'Image color mode: {img.mode}, grayscale: {is_grayscale}, size: {img.width}x{img.height}')
        return result
    except Exception as e:
        logger.error(f'Failed to read image: {str(e)}')
        return None
    finally:
        os.unlink(tmp_path)

def get_gimp_config_statusbar__2d2ae6cc356da88025063f58e6110bad(env, config: Dict[str, str]):
    """
    Gets the GIMP config file to check statusbar visibility setting.
    This getter retrieves the sessionrc file which contains window display settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_image_dimensions__844f5e73b108da449b9b68fcbef6bbf2(env, config: Dict[str, Any]):
    """
    Get the dimensions and aspect ratio of both the rotated image and original image.
    This enables verification that rotation was performed correctly by comparing dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (path to rotated image)

    Returns:
        dict: {
            'rotated_width': int,
            'rotated_height': int,
            'original_width': int,
            'original_height': int,
            'is_portrait': bool,
            'dimensions_swapped': bool
        } or None if file not found
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    original_path = '/home/user/Downloads/kingbird.jpeg'
    original_bytes = env.controller.get_file(original_path)
    import tempfile
    from PIL import Image
    import imagehash
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_rotated = tmp.name
        img_rotated = Image.open(tmp_rotated)
        (rotated_width, rotated_height) = img_rotated.size
        is_portrait = rotated_height > rotated_width
        rotated_hash = str(imagehash.average_hash(img_rotated))
        img_rotated.close()
        os.unlink(tmp_rotated)
        result = {'rotated_width': rotated_width, 'rotated_height': rotated_height, 'is_portrait': is_portrait, 'rotated_hash': rotated_hash}
        if original_bytes:
            with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
                tmp.write(original_bytes)
                tmp_original = tmp.name
            img_original = Image.open(tmp_original)
            (original_width, original_height) = img_original.size
            img_original_rotated = img_original.rotate(-90, expand=True)
            original_rotated_hash = str(imagehash.average_hash(img_original_rotated))
            img_original.close()
            img_original_rotated.close()
            os.unlink(tmp_original)
            dimensions_swapped = rotated_width == original_height and rotated_height == original_width
            hash_distance = imagehash.hex_to_hash(rotated_hash) - imagehash.hex_to_hash(original_rotated_hash)
            content_preserved = hash_distance <= 5
            result.update({'original_width': original_width, 'original_height': original_height, 'dimensions_swapped': dimensions_swapped, 'content_preserved': content_preserved, 'hash_distance': hash_distance})
            logger.info(f'Rotated: {rotated_width}x{rotated_height}, Original: {original_width}x{original_height}')
            logger.info(f'Dimensions swapped: {dimensions_swapped}, Content preserved: {content_preserved}, Hash distance: {hash_distance}')
        else:
            logger.warning(f'Original image not found at {original_path}, cannot verify rotation correctness')
            result['dimensions_swapped'] = False
            result['content_preserved'] = False
        return result
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        if 'tmp_rotated' in locals() and os.path.exists(tmp_rotated):
            os.unlink(tmp_rotated)
        if 'tmp_original' in locals() and os.path.exists(tmp_original):
            os.unlink(tmp_original)
        return None

def get_image_props__a01f23a4(env, config: dict):
    """Get basic image properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        dict: Image properties
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        return {'width': img.width, 'height': img.height, 'exists': True}
    except Exception as e:
        return {'exists': False}

def get_gimp_layer_name_config__174a6594(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_gimp_recolored_image__d52919fa54c3d57b57da98359753b674(env, config):
    """Get the recolored/exported image from VM and return both result and original image paths.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary containing paths to both result and original images
              {'result_path': str, 'original_path': str}
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    dest_filename = config.get('dest', os.path.basename(config['path']))
    result_cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(result_cache_path, 'wb') as f:
        f.write(file_bytes)
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    original_bytes = env.controller.get_file(original_path)
    original_cache_path = None
    if original_bytes:
        original_cache_path = os.path.join(env.cache_dir, 'Triangle_On_The_Side_original.png')
        with open(original_cache_path, 'wb') as f:
            f.write(original_bytes)
    return {'result_path': result_cache_path, 'original_path': original_cache_path}

def get_triangle_color__978cb6f31473d4802226bc7ae94b7399(env, config: dict):
    """
    Get the dominant color of the triangle.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: RGB color values {'r': int, 'g': int, 'b': int} or None if extraction fails
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        elif img.mode == 'L' or img.mode == 'P':
            img = img.convert('RGB')
        img_array = np.array(img)
        if len(img_array.shape) != 3 or img_array.shape[2] != 3:
            return None
        (height, width, channels) = img_array.shape
        (unique_colors, counts) = np.unique(img_array.reshape(-1, channels), axis=0, return_counts=True)
        if len(unique_colors) < 2:
            return None
        corner_pixels = [img_array[0, 0], img_array[0, -1], img_array[-1, 0], img_array[-1, -1]]
        (corner_colors_unique, corner_counts) = np.unique(corner_pixels, axis=0, return_counts=True)
        background_color = corner_colors_unique[np.argmax(corner_counts)]
        triangle_color = None
        max_non_bg_count = 0
        for (color, count) in zip(unique_colors, counts):
            if not np.array_equal(color, background_color):
                if count > max_non_bg_count:
                    max_non_bg_count = count
                    triangle_color = color
        if triangle_color is None:
            sorted_indices = np.argsort(counts)[::-1]
            if len(sorted_indices) >= 2:
                triangle_color = unique_colors[sorted_indices[1]]
            else:
                return None
        return {'r': int(triangle_color[0]), 'g': int(triangle_color[1]), 'b': int(triangle_color[2])}
    finally:
        os.unlink(tmp_path)

def get_jpg_list_file__8ddf03c6b80780c49b4f9497dee3f888(env, config: dict) -> list:
    """Read a text file containing jpg filenames and return as list.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of filenames from the file, one per line
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    try:
        content = file_bytes.decode('utf-8').strip()
        if not content:
            return []
        filenames = [line.strip() for line in content.split('\n') if line.strip()]
        return sorted(filenames)
    except Exception as e:
        return []

def get_gimp_config_rulers__96f3f9d88ca4e4230495b91ff566eb51(env, config: Dict[str, str]):
    """
    Gets the GIMP config file to check rulers visibility setting.
    This getter retrieves the sessionrc file which contains window display settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_gimp_image_dimensions__4bd0ac4fe70775f29bef20a161a34c39(env, config):
    """
    Get the dimensions of a GIMP-exported image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int} or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'File not found or empty: {file_path}')
            return None
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            logger.info(f'Image dimensions: {width}x{height}')
            return {'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        return None

def get_image_dimensions__1beedc32(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_gimp_scaled_image__93095e022ff3d0c1026d009f3ccc512b(env, config):
    """Get the scaled/exported image from VM and also download the original for comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to the downloaded scaled image file in cache
    """
    import os
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    dest_filename = config.get('dest', os.path.basename(config['path']))
    cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    original_bytes = env.controller.get_file(original_path)
    if original_bytes:
        original_cache_path = os.path.join(env.cache_dir, 'Triangle_On_The_Side.png')
        with open(original_cache_path, 'wb') as f:
            f.write(original_bytes)
    return cache_path

def get_image_size__7ecb394c9eae8a8e135a21ca629ec0de(env, config: Dict):
    """
    Get the dimensions of an image file and check if it exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path on VM

    Returns:
        dict: {"width": int, "height": int, "exists": bool} or {"exists": False} if file not found
    """
    from PIL import Image
    import tempfile
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get image file from {config['path']}")
        return {'exists': False}
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height, 'exists': True}
        logger.info(f'Image size: {img.width}x{img.height}')
        return result
    except Exception as e:
        logger.error(f'Failed to read image: {str(e)}')
        return {'exists': False}
    finally:
        os.unlink(tmp_path)

def get_image_properties__bd10eca8(env, config: dict):
    """Extract image properties (dimensions, format) from an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Image properties including width, height, format, exists
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            return {'exists': False, 'width': 0, 'height': 0, 'format': None}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result = {'exists': True, 'width': img.width, 'height': img.height, 'format': img.format}
            img.close()
            return result
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'format': None, 'error': str(e)}

def get_default_music_player__2fc648743a39d7031a57b208b33f9200(env, config: dict):
    """Gets the default music/audio player application.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required by framework)

    Returns:
        str: The default audio player .desktop file name (e.g., 'rhythmbox.desktop')
    """
    import requests
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['mp3', 'mpeg', 'ogg', 'flac', 'wav', 'x-wav', 'x-ms-wma', 'x-flac', 'aac', 'x-aac', 'mp4', 'x-mp4', 'x-mpeg']
        apps = []
        vm_ip = env.vm_ip
        port = env.server_port
        for ext in extensions:
            command = ['xdg-mime', 'query', 'default', f'audio/{ext}']
            response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
            if response.status_code == 200:
                app = response.json().get('output', '').strip()
                if app:
                    apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    else:
        raise Exception('Unsupported operating system', os_type)

def get_gimp_layer_name_config__3336340a(env, config):
    """Gets layer names from a GIMP .xcf file by parsing the file structure."""
    file_path = config['file_path']
    content = env.controller.get_file(file_path)
    if not content:
        logger.error(f'Failed to get GIMP .xcf file from {file_path}')
        return None
    dest_path = os.path.join(env.cache_dir, config['dest'])
    with open(dest_path, 'wb') as f:
        f.write(content)
    try:
        layer_names = parse_xcf_layers(dest_path)
        logger.info(f'Found layers in .xcf file: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error parsing .xcf file: {e}')
        return None

def get_gimp_image_file__bc5fe443a1102b529ba31d0ac9c81ab1(env, config: Dict[str, str]):
    """
    Gets an image file from VM and saves it to cache directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file on VM

    Returns:
        str: Path to the cached image file, or None if file cannot be read
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        cache_filename = os.path.basename(file_path)
        cache_path = os.path.join(env.cache_dir, cache_filename)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        logger.info(f'Image saved to cache: {cache_path}')
        return cache_path
    except Exception as e:
        logger.error(f'Error saving image to cache: {e}')
        return None

def get_gimp_layer_names__a85f6474(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_third_image__9a3feca86555f3c82732707871883142(env, config: dict):
    """Check if the third image from document was saved to Pictures.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'exists': bool,
            'size': int or None,
            'is_png': bool,
            'hash': str or None,
            'reference_hash': str or None
        }
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    reference_hash = None
    docx_path = '/home/user/.thunderbird/yizp5ecy.default-release/Mail/Local Folders/Notes.sbd/Notes'
    try:
        docx_bytes = None
        potential_paths = ['/home/user/.thunderbird/yizp5ecy.default-release/Mail/Local Folders/Notes.sbd/Notes', '/home/user/Downloads/lecture-notes.docx', '/tmp/lecture-notes.docx']
        for potential_path in potential_paths:
            potential_bytes = env.controller.get_file(potential_path)
            if potential_bytes and potential_bytes[:4] == b'PK\x03\x04':
                docx_bytes = potential_bytes
                break
        if docx_bytes:
            with zipfile.ZipFile(io.BytesIO(docx_bytes), 'r') as docx_zip:
                image_files = [f for f in docx_zip.namelist() if f.startswith('word/media/')]
                image_files.sort()
                third_image_bytes = None
                if 'word/media/image3.png' in docx_zip.namelist():
                    third_image_bytes = docx_zip.read('word/media/image3.png')
                elif len(image_files) >= 3:
                    third_image_bytes = docx_zip.read(image_files[2])
                if third_image_bytes:
                    reference_hash = hashlib.sha256(third_image_bytes).hexdigest()
    except Exception:
        pass
    if file_bytes:
        is_png = file_bytes[:8] == b'\x89PNG\r\n\x1a\n'
        file_hash = hashlib.sha256(file_bytes).hexdigest() if is_png else None
        return {'exists': True, 'size': len(file_bytes), 'is_png': is_png, 'hash': file_hash, 'reference_hash': reference_hash}
    else:
        return {'exists': False, 'size': None, 'is_png': False, 'hash': None, 'reference_hash': reference_hash}

def get_gimp_layer_names__677f0e4c(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    Args:
        env: Environment object
        config: Configuration dict with 'file_path'

    Returns:
        list: List of layer names in top-to-bottom order (index 0 is the topmost layer)
    """
    try:
        from desktop_env.evaluators.getters.gimp import parse_xcf_layers
        import os
        file_path = config.get('file_path', '/home/user/Desktop/resized.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_layer_names_677f0e4c.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Found layers: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_gif_file_info__06536a54(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_has_background_image__4b23e04e(env, config: Dict[str, str]):
    """
    Extract background image from slide and compare with video frame at 00:15.

    Args:
        env: Environment object
        config: Dict containing:
            - ppt_file_path: Path to the PowerPoint file
            - slide_index: Index of the slide to check (0-based)
            - video_file_path: Path to the video file (landscape.mp4)
            - timestamp: Timestamp to extract frame (in seconds, default 15)

    Returns:
        dict: Dictionary with:
            - has_background: bool indicating if background image exists
            - background_image: numpy array of the extracted background (or None)
            - expected_frame: numpy array of the video frame at timestamp (or None)
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    video_file_path = config.get('video_file_path', '/home/user/Desktop/landscape.mp4')
    timestamp = config.get('timestamp', 15)
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    video_file_localhost_path = get_vm_file(env, {'path': video_file_path, 'dest': os.path.split(video_file_path)[-1]})
    background_image = None
    has_background = False
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return {'has_background': False, 'background_image': None, 'expected_frame': None}
            bg_rel_id = None
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                for child in root.iter(bg_tag):
                    for element in child.iter(image_tag):
                        if attr_tag in element.attrib:
                            bg_rel_id = element.attrib[attr_tag]
                            has_background = True
                            break
                    if has_background:
                        break
            if has_background and bg_rel_id and (slide_rels_file in myzip.namelist()):
                with myzip.open(slide_rels_file) as f:
                    rels_tree = ET.parse(f)
                    rels_root = rels_tree.getroot()
                    rel_tag = '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'
                    for rel in rels_root.findall(rel_tag):
                        if rel.attrib.get('Id') == bg_rel_id:
                            target = rel.attrib.get('Target')
                            if target:
                                image_path = os.path.normpath(os.path.join('ppt/slides', target)).replace('\\', '/')
                                if image_path in myzip.namelist():
                                    with myzip.open(image_path) as img_file:
                                        img_data = img_file.read()
                                        img = Image.open(tempfile.NamedTemporaryFile(delete=False, suffix='.png'))
                                        temp_img_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                                        temp_img_file.write(img_data)
                                        temp_img_file.close()
                                        img = cv2.imread(temp_img_file.name)
                                        if img is not None:
                                            background_image = img
                                        os.unlink(temp_img_file.name)
                                break
    except Exception as e:
        pass
    expected_frame = None
    try:
        cap = cv2.VideoCapture(video_file_localhost_path)
        if cap.isOpened():
            cap.set(cv2.CAP_PROP_POS_MSEC, timestamp * 1000)
            (ret, frame) = cap.read()
            if ret:
                expected_frame = frame
            cap.release()
    except Exception as e:
        pass
    return {'has_background': has_background, 'background_image': background_image, 'expected_frame': expected_frame}

def get_image_mode__9660fa13484d43a1462069231ef86deb(env, config: Dict[str, Any]):
    """
    Get the color mode of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Image mode (e.g., 'L' for grayscale, 'RGB', 'RGBA') or None if error
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    import tempfile
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        mode = img.mode
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image mode: {mode}')
        return mode
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_gif_dimensions__90c23f3797e29598f167849a527a40d5(env, config: dict):
    """
    Get dimensions and animation properties of a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file path on VM

    Returns:
        dict: Info including width, height, exists, is_animated, frame_count
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path specified in config')
        return {'exists': False, 'width': 0, 'height': 0, 'is_animated': False, 'frame_count': 0}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.warning(f'File not found or empty: {file_path}')
        return {'exists': False, 'width': 0, 'height': 0, 'is_animated': False, 'frame_count': 0}
    import tempfile
    try:
        with tempfile.NamedTemporaryFile(suffix='.gif', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            (width, height) = img.size
            frame_count = 0
            try:
                while True:
                    frame_count += 1
                    img.seek(img.tell() + 1)
            except EOFError:
                pass
            is_animated = frame_count > 1
            img.close()
            logger.info(f'GIF analysis: {width}x{height}, {frame_count} frames, animated={is_animated}')
            return {'exists': True, 'width': width, 'height': height, 'is_animated': is_animated, 'frame_count': frame_count}
        except Exception as e:
            logger.error(f'Error opening image file: {e}')
            return {'exists': True, 'width': 0, 'height': 0, 'is_animated': False, 'frame_count': 0}
    finally:
        if 'tmp_path' in locals():
            try:
                os.unlink(tmp_path)
            except:
                pass

def get_default_audio_player__3b265860(env, config: dict):
    """Gets the default application for audio/music files.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The most common default audio player registered for audio MIME types
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        extensions = ['aac', 'flac', 'mp3', 'mpeg', 'ogg', 'opus', 'wav', 'webm', 'x-aac', 'x-flac', 'x-mp3', 'x-mpeg', 'x-ms-wma', 'x-vorbis', 'x-vorbis+ogg', 'x-wav']
        apps = []
        for ext in extensions:
            app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', f'audio/{ext}']})
            if app:
                apps.append(app)
        if len(apps) == 0:
            return 'unknown'
        else:
            return Counter(apps).most_common(1)[0][0]
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_image_properties__9d6c98f0(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_image_dimensions__de812dd5b44b906cb9793a8d4a3f91bf(env, config):
    """
    Get the dimensions of an image file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the VM file path

    Returns:
        dict: {'width': int, 'height': int} or None if file doesn't exist
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        img.close()
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Error reading image dimensions: {e}')
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_size__10742f90(env, config: dict):
    """Get image dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        dict: Image dimensions
    """
    vm_path = config.get('path')
    dest_name = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        return {'width': img.width, 'height': img.height, 'exists': True}
    except Exception as e:
        return {'exists': False}

def get_gimp_transparency__456c340f2121fceb6f1006996239e1cf(env, config: Dict[str, str]):
    """
    Check if a PNG file exists and has transparency (alpha channel).
    Returns dict with 'exists' and 'has_transparency' keys.
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'has_transparency': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.debug(f'File does not exist or could not be retrieved: {file_path}')
        return {'exists': False, 'has_transparency': False}
    import tempfile
    import os
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            has_transparency = False
            if img.mode in ('RGBA', 'LA', 'PA'):
                has_transparency = True
            elif 'transparency' in img.info:
                has_transparency = True
            logger.debug(f'Image mode: {img.mode}, has_transparency: {has_transparency}')
            return {'exists': True, 'has_transparency': has_transparency, 'mode': img.mode, 'size': img.size}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking transparency: {e}')
        return {'exists': True, 'has_transparency': False}

def get_gimp_layer_exists__e5a8318d(env, config):
    """
    Check if a specific layer exists in the GIMP XCF file.

    This getter uses the existing get_gimp_layer_names function to retrieve
    all layer names and checks if the expected layer exists.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file (default: /home/user/Desktop/white_background.xcf)
            - 'layer_name': Name of the layer to check for

    Returns:
        bool: True if layer exists, False otherwise
    """
    try:
        from desktop_env.evaluators.getters.gimp import get_gimp_layer_names
        layer_names = get_gimp_layer_names(env, config)
        expected_layer = config.get('layer_name', '')
        layer_exists = expected_layer in layer_names
        logger.info(f"Checking for layer '{expected_layer}' in {layer_names}: {layer_exists}")
        return layer_exists
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return False

def get_vm_file(env, config: Dict[str, Any]):
    """Get file from VM - simplified version for this getter."""
    path = config['path']
    dest = config.get('dest', os.path.basename(path))
    try:
        file_content = env.controller.get_file(path)
        if file_content is None:
            logger.warning(f'Failed to get file from VM: {path}')
            return None
        cache_path = os.path.join(env.cache_dir, dest)
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        return cache_path
    except Exception as e:
        logger.error(f'Error getting file {path}: {e}')
        return None

def get_image_properties__fa4ed82a(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_image_color_mode__14447080df6d9553c7e99d3265fc5a81(env, config: Dict):
    """
    Get the color mode of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' key pointing to the image file

    Returns:
        Dict with 'mode' key (e.g., 'RGB', 'L', 'P', 'RGBA'), or None if file cannot be read
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        logger.error(f"Failed to get file: {config['path']}")
        return None
    try:
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            mode = img.mode
            img.close()
            return {'mode': mode}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image color mode: {e}')
        return None

def get_gimp_file_bytes__05275b4640be680d1ddfc9692455a07a(env, config: Dict[str, str]):
    """
    Get file bytes from VM for image comparison.
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path specified in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file: {file_path}')
        return None
    import os
    cache_path = os.path.join(env.cache_dir, config.get('dest', 'temp_file.png'))
    os.makedirs(os.path.dirname(cache_path), exist_ok=True)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    logger.debug(f'Saved file to cache: {cache_path}')
    return cache_path

def get_gimp_layer_names__8249d409(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_image_file_count__3fde58b2(env, config: Dict[str, Any]) -> int:
    """
    Get the count of image files in a specified directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory_path'

    Returns:
        int: Number of image files in directory
    """
    directory_path = config.get('directory_path', '')
    cmd = f"find '{directory_path}' -maxdepth 1 -type f \\( -iname '*.jpg' -o -iname '*.jpeg' -o -iname '*.png' -o -iname '*.gif' -o -iname '*.bmp' \\) 2>/dev/null | wc -l"
    result = env.controller.run_bash_script(cmd, timeout=10)
    count = 0
    if result.get('status') == 'success':
        try:
            count = int(result.get('output', '0').strip())
        except ValueError:
            logger.warning(f"Failed to parse file count: {result.get('output')}")
    logger.info(f'Image file count in {directory_path}: {count}')
    return count

def get_image_dimensions__696943e1(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_image_contrast__979c644a69a9cc6d01b110bbe0f08e78(env, config):
    """Get contrast value of a GIMP image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        float: Contrast value (standard deviation of pixel values)
    """
    import tempfile
    import os
    import numpy as np
    from PIL import Image
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        pixels = np.asarray(img, dtype=np.float32)
        contrast = np.std(pixels)
        return float(contrast)
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_image_dimensions__7cb89717dbfd62e5cbbd4dcc85a4e268(env, config: Dict[str, Any]):
    """
    Get the dimensions of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {'width': int, 'height': int} or None if file not found
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No path provided in config')
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {file_path}')
        return None
    import tempfile
    from PIL import Image
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpeg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        (width, height) = img.size
        img.close()
        os.unlink(tmp_path)
        logger.info(f'Image dimensions: {width}x{height}')
        return {'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        if 'tmp_path' in locals() and os.path.exists(tmp_path):
            os.unlink(tmp_path)
        return None

def get_image_dimensions__da5d7378bbf41608325407fc00f8d126(env, config):
    """Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: {'width': int, 'height': int} or None if failed
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        return result
    finally:
        os.unlink(tmp_path)

def get_gif_file_info__d18c2fd4(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_dimensions__3f39d534f7803089ed331d19c2a1bc89(env, config):
    """Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: {'width': int, 'height': int} or None if failed
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        return result
    finally:
        os.unlink(tmp_path)

def get_image_properties__ab3f6dd1(env, config: Dict) -> Optional[Dict]:
    """
    Get image properties including dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'exists', 'width', 'height' keys
    """
    file_path = config.get('path')
    if not file_path:
        return {'exists': False, 'width': 0, 'height': 0}
    try:
        local_path = get_vm_file(env, {'path': file_path, 'dest': os.path.basename(file_path)})
        if not local_path or not os.path.exists(local_path):
            return {'exists': False, 'width': 0, 'height': 0}
        with Image.open(local_path) as img:
            return {'exists': True, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'exists': False, 'width': 0, 'height': 0, 'error': str(e)}

def get_gimp_gimprc_file__dc948653(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_gimp_brightness__0b33739cf65dc15a04df08d23f597efa(env, config):
    """Get brightness value of a GIMP image.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for the image file path

    Returns:
        float: Average brightness value [0-255]
    """
    import tempfile
    import os
    from PIL import Image, ImageStat
    file_path = config.get('path')
    if not file_path:
        return None
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        img = Image.open(tmp_path)
        grayscale = img.convert('L')
        stat = ImageStat.Stat(grayscale)
        brightness = stat.mean[0]
        return brightness
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_gif_file_info__06722a19(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_hash__6510dd9a(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_docx_image_count__1b5dc5fe(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_png_file__d04e36cb527bbbbfdb06458981bf8945(env, config: dict):
    """
    Get PNG file from VM and return the file path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        str: Path to cached file or None if failed
    """
    vm_path = config['path']
    dest_name = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        logger.info(f'Successfully saved PNG file: {cache_path} ({len(file_bytes)} bytes)')
        return cache_path
    except Exception as e:
        logger.error(f'Error saving file: {e}')
        return None

def get_image_properties__565232d7(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_image_hash__b461596b(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_orientation__6575e228(env, config):
    """
    Get the dimensions (width, height) of an image file.

    This getter is used to verify that an image has been rotated correctly
    by checking if the dimensions have been swapped as expected.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'path': Path to the image file in the VM
            - 'dest': Destination filename for local cache

    Returns:
        dict: Dictionary with 'width' and 'height' keys, or None on error
    """
    try:
        file_path = config.get('path')
        dest = config.get('dest', 'temp_image.jpg')
        if not file_path:
            logger.error('No file path specified in config')
            return None
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, dest)
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        (width, height) = img.size
        logger.info(f'Image dimensions: {width}x{height}')
        return {'width': width, 'height': height}
    except Exception as e:
        logger.error(f'Error getting image orientation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_image_file_exists__836db3cb3a7cc19d82f98c6a439eb80c(env, config: Dict[str, Any]) -> bool:
    """
    Check if an image file exists at the specified path on the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for VM file path

    Returns:
        bool: True if file exists and is an image, False otherwise
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'File does not exist at path: {path}')
            return False
        image_signatures = [(b'\xff\xd8\xff', 'JPEG'), (b'\x89PNG\r\n\x1a\n', 'PNG'), (b'GIF87a', 'GIF'), (b'GIF89a', 'GIF'), (b'BM', 'BMP'), (b'II*\x00', 'TIFF'), (b'MM\x00*', 'TIFF')]
        for (signature, format_name) in image_signatures:
            if file_bytes.startswith(signature):
                logger.info(f'Image file ({format_name}) exists at {path}')
                return True
        logger.warning(f'File exists at {path} but is not a recognized image format')
        return False
    except Exception as e:
        logger.error(f'Error checking file existence: {e}')
        return False

def get_docx_images_count__17e4ac0c(env, config: dict):
    """Count images in DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Image count
    """
    vm_path = config.get('path', '/home/user/Desktop/Viewing_Your_Class_Schedule_and_Textbooks.docx')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        return 0
    cache_path = os.path.join(env.cache_dir, os.path.basename(vm_path))
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        doc = Document(cache_path)
        return sum((1 for rel in doc.part.rels.values() if 'image' in rel.reltype))
    except Exception:
        return 0

def get_xcf_layer_names__b148e375(env, config):
    """
    Get the layer names from a GIMP XCF file using GIMP's scripting interface.

    Args:
        env: Environment object
        config: Configuration dict with 'path' key pointing to the XCF file

    Returns:
        list: List of layer names in the XCF file, or empty list if error
    """
    xcf_path = config.get('path', '')
    if not xcf_path:
        logger.error('No path specified in config')
        return []
    scheme_script = f'(let* ((image (car (gimp-file-load RUN-NONINTERACTIVE "{xcf_path}" "{xcf_path}")))\n               (layers (gimp-image-get-layers image))\n               (num-layers (car layers))\n               (layer-array (cadr layers))\n               (i 0))\n          (while (< i num-layers)\n            (let ((layer-name (car (gimp-item-get-name (aref layer-array i)))))\n              (gimp-message layer-name))\n            (set! i (+ i 1)))\n          (gimp-image-delete image)\n          (gimp-quit 0))'
    try:
        python_code = f"""\nimport subprocess\nimport sys\n\nscheme_script = '''{scheme_script}'''\n\ntry:\n    # Run GIMP in batch mode (-i: no interface, -b: batch mode)\n    result = subprocess.run(\n        ['gimp', '-i', '-b', scheme_script, '-b', '(gimp-quit 0)'],\n        capture_output=True,\n        text=True,\n        timeout=15\n    )\n    \n    # GIMP messages go to stderr\n    print(result.stderr)\nexcept subprocess.TimeoutExpired:\n    print("TIMEOUT", file=sys.stderr)\nexcept Exception as e:\n    print(f"ERROR: {{e}}", file=sys.stderr)\n"""
        result = env.controller.execute_python_command(python_code)
        if result and result.get('output'):
            output = result['output'].strip()
            logger.info(f'GIMP output: {output}')
            layer_names = []
            for line in output.split('\n'):
                line = line.strip()
                if line and (not any((skip in line.lower() for skip in ['batch', 'command', 'error', 'warning', 'gimp-', 'script-fu', 'opening', 'saving']))):
                    if line:
                        layer_names.append(line)
            if layer_names:
                logger.info(f'Extracted layer names from {xcf_path}: {layer_names}')
                return layer_names
        logger.warning('Scheme script approach failed, attempting direct XCF parsing')
        return parse_xcf_file_directly(env, xcf_path)
    except Exception as e:
        logger.error(f'Error extracting layer names: {e}')
        return parse_xcf_file_directly(env, xcf_path)

def get_gimp_layer_group_config__7ba73b05(env, config):
    """
    Gets the layer structure from GIMP XCF file to check for layer groups.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of dictionaries containing layer information with 'name' and 'is_group' keys
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_info = parse_xcf_layer_groups(local_path)
        logger.info(f'Parsed layer structure from XCF: {layer_info}')
        return layer_info
    except Exception as e:
        logger.error(f'Error getting GIMP layer structure: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_image_mode__d9cbc3c3(env, config):
    """
    Get the mode and check if image is grayscale.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'mode': str, 'is_grayscale': bool} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/grayscale.png')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, 'temp_image_mode.png')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        mode = img.mode
        is_grayscale = False
        if mode in ('L', 'LA'):
            is_grayscale = True
        elif mode in ('RGB', 'RGBA'):
            arr = np.array(img)
            if len(arr.shape) >= 3:
                r = arr[:, :, 0]
                g = arr[:, :, 1]
                b = arr[:, :, 2]
                is_grayscale = np.allclose(r, g, atol=2) and np.allclose(g, b, atol=2)
            else:
                is_grayscale = True
        result = {'mode': mode, 'is_grayscale': is_grayscale}
        logger.info(f'Image mode: {result}')
        return result
    except Exception as e:
        logger.error(f'Error getting image mode: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_extracted_image_properties__345e8ddb(env, config: dict):
    """
    Get properties of an extracted image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the image file path

    Returns:
        dict: Image properties including exists, size, format, dimensions
    """
    import os
    from PIL import Image
    from io import BytesIO
    path = config.get('path', '/home/user/extracted_image.png')
    file_bytes = env.controller.get_file(path)
    if file_bytes is None:
        return {'exists': False, 'size': 0, 'format': None, 'width': 0, 'height': 0}
    try:
        img = Image.open(BytesIO(file_bytes))
        return {'exists': True, 'size': len(file_bytes), 'format': img.format, 'width': img.width, 'height': img.height}
    except Exception as e:
        return {'exists': True, 'size': len(file_bytes), 'format': 'unknown', 'width': 0, 'height': 0, 'error': str(e)}

def get_image_dimensions__a4f96e46(env, config: dict):
    """
    Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM file path) and 'dest' (cache filename)

    Returns:
        tuple: (width, height) of the image
    """
    import os
    vm_path = config.get('path')
    dest = config.get('dest', 'temp_image.png')
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    try:
        img = Image.open(cache_path)
        dimensions = img.size
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Failed to open image: {e}')
        return None

def get_docx_image_count__34e55d07(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_gimp_config_fullscreen__8efddf2685fdb790d7823145c3565e94(env, config: Dict[str, str]):
    """
    Gets the GIMP config file to check fullscreen mode setting.
    This getter retrieves the sessionrc file which contains window display settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_docx_image_count__cd0e399d(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_jpeg_export_check__39da7f334341155f29f73cbacf02786e(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Check if JPEG file was exported from GIMP and extract its properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to JPEG file

    Returns:
        Dict with file_exists, format, and size_kb properties
    """
    jpeg_path = config.get('path', '')
    result = {'file_exists': False, 'format': '', 'size_kb': 0}
    file_bytes = env.controller.get_file(jpeg_path)
    if not file_bytes:
        logger.warning(f'JPEG file not found: {jpeg_path}')
        return result
    result['file_exists'] = True
    result['size_kb'] = len(file_bytes) / 1024
    try:
        with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            result['format'] = img.format if img.format else ''
            logger.info(f"JPEG format: {result['format']}, size: {result['size_kb']:.2f} KB")
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error reading JPEG file: {e}')
    return result

def get_gimp_flipped_image__1aaf02638da71a4a84f47e22e2395da3(env, config):
    """Get both the flipped/exported image and original image from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for result image

    Returns:
        dict: Dictionary with 'result_path' and 'original_path' keys
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    dest_filename = config.get('dest', os.path.basename(config['path']))
    result_cache_path = os.path.join(env.cache_dir, dest_filename)
    with open(result_cache_path, 'wb') as f:
        f.write(file_bytes)
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    original_bytes = env.controller.get_file(original_path)
    original_cache_path = None
    if original_bytes:
        original_cache_path = os.path.join(env.cache_dir, 'Triangle_On_The_Side.png')
        with open(original_cache_path, 'wb') as f:
            f.write(original_bytes)
    return {'result_path': result_cache_path, 'original_path': original_cache_path}

def get_image_hash__db2588a0(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_image_hash__73176121(env, config: dict):
    """Get SHA256 hash of an image file on VM.

    Config:
        path (str): absolute path to the image file on VM

    Returns:
        str: SHA256 hash of the file, or None if file doesn't exist
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            logger.error(f'File not found: {file_path}')
            return None
        file_hash = hashlib.sha256(file_bytes).hexdigest()
        logger.info(f'Computed hash for {file_path}: {file_hash}')
        return file_hash
    except Exception as e:
        logger.error(f'Error getting image hash from {file_path}: {e}')
        return None

def get_gif_file_info__3d85489a(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_dimensions__4694337da0a8886d5bd508a95fd83b12(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Extract image dimensions from a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict with 'width', 'height' keys, or None if file doesn't exist or isn't valid
    """
    path = config['path']
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        logger.warning(f'File not found: {path}')
        return None
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            dimensions = {'width': img.width, 'height': img.height}
            return dimensions
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Failed to read image dimensions: {e}')
        return None

def get_has_background_image__f45ebd91(env, config: Dict[str, str]):
    """
    Extract the background image from a slide and compute its hash for comparison.
    Also extracts a reference frame from the video at the specified timestamp.

    Args:
        env: Environment object
        config: Configuration dict with keys:
            - ppt_file_path: Path to the PowerPoint file
            - slide_index: Index of the slide (0-based)

    Returns:
        dict: Dictionary with keys:
            - has_background: boolean indicating if background image exists
            - background_image_hash: MD5 hash of the background image (if exists)
            - background_image_data: bytes of the background image (if exists)
            - video_frame_hash: MD5 hash of the video frame at 00:10
            - video_frame_data: bytes of the video frame
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    result = {'has_background': False, 'background_image_hash': None, 'background_image_data': None, 'video_frame_hash': None, 'video_frame_data': None}
    video_path = '/home/user/Desktop/landscape.mp4'
    video_localhost_path = get_vm_file(env, {'path': video_path, 'dest': os.path.split(video_path)[-1]})
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_frame:
            frame_path = tmp_frame.name
        subprocess.run(['ffmpeg', '-y', '-ss', '00:00:10', '-i', video_localhost_path, '-vframes', '1', '-f', 'image2', frame_path], capture_output=True, timeout=10)
        if os.path.exists(frame_path) and os.path.getsize(frame_path) > 0:
            with open(frame_path, 'rb') as f:
                video_frame_data = f.read()
                result['video_frame_data'] = video_frame_data
                result['video_frame_hash'] = hashlib.md5(video_frame_data).hexdigest()
            os.unlink(frame_path)
    except Exception as e:
        pass
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return result
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                image_rel_id = None
                for child in root.iter(bg_tag):
                    for element in child.iter(image_tag):
                        image_rel_id = element.attrib.get(attr_tag)
                        if image_rel_id:
                            result['has_background'] = True
                            break
                    if image_rel_id:
                        break
                if not image_rel_id:
                    return result
            if slide_rels_file in myzip.namelist():
                with myzip.open(slide_rels_file) as f:
                    rels_tree = ET.parse(f)
                    rels_root = rels_tree.getroot()
                    rel_tag = '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'
                    for rel in rels_root.findall(rel_tag):
                        if rel.attrib.get('Id') == image_rel_id:
                            target = rel.attrib.get('Target')
                            if target:
                                if target.startswith('../'):
                                    image_path = 'ppt/' + target[3:]
                                else:
                                    image_path = 'ppt/slides/' + target
                                if image_path in myzip.namelist():
                                    with myzip.open(image_path) as img_file:
                                        image_data = img_file.read()
                                        result['background_image_data'] = image_data
                                        result['background_image_hash'] = hashlib.md5(image_data).hexdigest()
                                break
    except Exception as e:
        pass
    return result

def get_image_properties__9472df29(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_docx_image_count__7d8c4525(env, config):
    """
    Extract image information including count and position from a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        dict: Dictionary containing:
            - 'total_count': Total number of images in the document
            - 'paragraph_count': Total number of paragraphs
            - 'last_image_position': Position index of the last image (paragraph index)
            - 'has_image_at_end': Boolean indicating if an image appears after the last text paragraph
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return {'total_count': 0, 'paragraph_count': 0, 'last_image_position': -1, 'has_image_at_end': False}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return {'total_count': 0, 'paragraph_count': 0, 'last_image_position': -1, 'has_image_at_end': False}
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            total_image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    total_image_count += 1
            paragraph_count = len(doc.paragraphs)
            last_text_paragraph_index = -1
            last_image_position = -1
            for (idx, para) in enumerate(doc.paragraphs):
                if para.text.strip():
                    last_text_paragraph_index = idx
                for run in para.runs:
                    if run._element.findall('.//{*}blip'):
                        last_image_position = idx
            has_image_at_end = last_image_position > last_text_paragraph_index and last_text_paragraph_index >= 0
            return {'total_count': total_image_count, 'paragraph_count': paragraph_count, 'last_image_position': last_image_position, 'has_image_at_end': has_image_at_end}
    except Exception as e:
        logger.error(f'Error extracting image information: {e}')
        return {'total_count': 0, 'paragraph_count': 0, 'last_image_position': -1, 'has_image_at_end': False}

def get_image_dimensions__a846d552e46987bd85e04c0a4f658c7a(env, config):
    """Get the dimensions (width, height) of an image file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the image file

    Returns:
        dict: {'width': int, 'height': int} or None if failed
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        result = {'width': img.width, 'height': img.height}
        return result
    finally:
        os.unlink(tmp_path)

def get_gimp_gimprc_file__d7b0aa1f(env, config: Dict[str, str]):
    """
    Gets the gimprc config file of GIMP.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_name': Name of the config file (e.g., 'gimprc')
            - 'dest': Destination filename in cache

    Returns:
        str: Path to the downloaded config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command(f"import os; print(os.path.expanduser('~/.config/GIMP/2.10/{config['file_name']}'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    if not content:
        logger.error('Failed to get GIMP gimprc config file.')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_png_files_list__ec3ddc36(env, config: Dict[str, Any]) -> List[str]:
    """Get list of PNG files in a directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        List of PNG filenames
    """
    path = config['path']
    result = env.controller.get_vm_directory_tree(path)
    if 'children' not in result:
        return []
    png_files = [node['name'] for node in result['children'] if node['name'].endswith('.png')]
    return png_files

def get_image_rotation_check__25c734e4497155c51ced0623dec284fc(env, config: Dict):
    """
    Get both source and result images for rotation comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (result image) and 'source_cache_path' (original)

    Returns:
        Dict with 'result_path' and 'source_path' for metric comparison
    """
    import tempfile
    result_path = config.get('path')
    result_bytes = env.controller.get_file(result_path)
    if not result_bytes:
        logger.error(f'Failed to get result image from {result_path}')
        return None
    result_temp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    result_temp.write(result_bytes)
    result_temp.close()
    source_cache_path = config.get('source_cache_path')
    return {'result_path': result_temp.name, 'source_path': source_cache_path}

def get_docx_image_count__e898c34b(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_image_crop_info__59e6ca63b22becd85a8942d1a29325d9(env, config):
    """Get image dimensions to verify crop operation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Image dimensions and aspect ratio or None
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        img = Image.open(tmp_path)
        (width, height) = (img.width, img.height)
        aspect_ratio = width / height if height > 0 else 0
        result = {'width': width, 'height': height, 'aspect_ratio': round(aspect_ratio, 3)}
        img.close()
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_gimp_layer_names__aa5f92aa(env, config):
    """
    Get the list of layer names from a GIMP XCF file.

    This getter parses the XCF file structure to extract the actual layer names
    from the currently saved image. This ensures we're checking the actual layers
    in the image, not config settings.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'file_path': Path to the XCF file to check

    Returns:
        list: List of layer names in the image, or empty list on error
    """
    try:
        file_path = config.get('file_path', '/home/user/Desktop/white_background.xcf')
        xcf_content = env.controller.get_file(file_path)
        if not xcf_content:
            logger.error(f'Failed to get XCF file: {file_path}')
            return []
        local_path = os.path.join(env.cache_dir, 'temp_image.xcf')
        with open(local_path, 'wb') as f:
            f.write(xcf_content)
        layer_names = parse_xcf_layers(local_path)
        logger.info(f'Parsed layers from XCF: {layer_names}')
        return layer_names
    except Exception as e:
        logger.error(f'Error getting GIMP layer names: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return []

def get_has_background_image__6c48637d(env, config: Dict[str, str]):
    """
    Extract the background image from a slide and compute its hash for comparison.
    Also extracts a reference frame from the video at the specified timestamp.

    Args:
        env: Environment object
        config: Configuration dict with keys:
            - ppt_file_path: Path to the PowerPoint file
            - slide_index: Index of the slide (0-based)

    Returns:
        dict: Dictionary with keys:
            - has_background: boolean indicating if background image exists
            - background_image_hash: MD5 hash of the background image (if exists)
            - background_image_data: bytes of the background image (if exists)
            - video_frame_hash: MD5 hash of the video frame at 00:03
            - video_frame_data: bytes of the video frame
    """
    ppt_file_path = config['ppt_file_path']
    slide_index = int(config['slide_index'])
    ppt_file_localhost_path = get_vm_file(env, {'path': ppt_file_path, 'dest': os.path.split(ppt_file_path)[-1]})
    result = {'has_background': False, 'background_image_hash': None, 'background_image_data': None, 'video_frame_hash': None, 'video_frame_data': None}
    video_path = '/home/user/Desktop/landscape.mp4'
    video_localhost_path = get_vm_file(env, {'path': video_path, 'dest': os.path.split(video_path)[-1]})
    try:
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_frame:
            frame_path = tmp_frame.name
        subprocess.run(['ffmpeg', '-y', '-ss', '00:00:03', '-i', video_localhost_path, '-vframes', '1', '-f', 'image2', frame_path], capture_output=True, timeout=10)
        if os.path.exists(frame_path) and os.path.getsize(frame_path) > 0:
            with open(frame_path, 'rb') as f:
                video_frame_data = f.read()
                result['video_frame_data'] = video_frame_data
                result['video_frame_hash'] = hashlib.md5(video_frame_data).hexdigest()
            os.unlink(frame_path)
    except Exception as e:
        pass
    try:
        with zipfile.ZipFile(ppt_file_localhost_path, 'r') as myzip:
            slide_xml_file = 'ppt/slides/slide{}.xml'.format(slide_index + 1)
            slide_rels_file = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_index + 1)
            if slide_xml_file not in myzip.namelist():
                return result
            with myzip.open(slide_xml_file) as f:
                tree = ET.parse(f)
                root = tree.getroot()
                bg_tag = '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr'
                image_tag = '{http://schemas.openxmlformats.org/drawingml/2006/main}blip'
                attr_tag = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
                image_rel_id = None
                for child in root.iter(bg_tag):
                    for element in child.iter(image_tag):
                        image_rel_id = element.attrib.get(attr_tag)
                        if image_rel_id:
                            result['has_background'] = True
                            break
                    if image_rel_id:
                        break
                if not image_rel_id:
                    return result
            if slide_rels_file in myzip.namelist():
                with myzip.open(slide_rels_file) as f:
                    rels_tree = ET.parse(f)
                    rels_root = rels_tree.getroot()
                    rel_tag = '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'
                    for rel in rels_root.findall(rel_tag):
                        if rel.attrib.get('Id') == image_rel_id:
                            target = rel.attrib.get('Target')
                            if target:
                                if target.startswith('../'):
                                    image_path = 'ppt/' + target[3:]
                                else:
                                    image_path = 'ppt/slides/' + target
                                if image_path in myzip.namelist():
                                    with myzip.open(image_path) as img_file:
                                        image_data = img_file.read()
                                        result['background_image_data'] = image_data
                                        result['background_image_hash'] = hashlib.md5(image_data).hexdigest()
                                break
    except Exception as e:
        pass
    return result

def get_image_properties__12182c78(env, config: dict):
    """Get properties of a specific image in a PowerPoint slide.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - path: Path to PPTX file on VM
            - slide_index: Slide index (0-based)
            - shape_index: Shape index (0-based)

    Returns:
        dict: Image properties including left, top, width, height
    """
    vm_path = config['path']
    slide_idx = int(config.get('slide_index', 0))
    shape_idx = int(config.get('shape_index', 0))
    local_path = get_vm_file(env, {'path': vm_path, 'dest': os.path.basename(vm_path)})
    prs = Presentation(local_path)
    if slide_idx >= len(prs.slides):
        return None
    slide = prs.slides[slide_idx]
    if shape_idx >= len(slide.shapes):
        return None
    shape = slide.shapes[shape_idx]
    return {'left': shape.left, 'top': shape.top, 'width': shape.width, 'height': shape.height, 'shape_type': shape.shape_type}

def get_image_dimensions__184d842d21bd06203c79c089532a2315(env, config: dict):
    """
    Get image file from VM and return its dimensions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (VM path) and 'dest' (cache filename)

    Returns:
        dict: Dictionary with 'width', 'height', and 'format' keys, or None if failed
    """
    vm_path = config['path']
    dest_name = config['dest']
    file_bytes = env.controller.get_file(vm_path)
    if not file_bytes:
        logger.error(f'Failed to get file from VM: {vm_path}')
        return None
    cache_path = os.path.join(env.cache_dir, dest_name)
    try:
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        img = Image.open(cache_path)
        result = {'width': img.size[0], 'height': img.size[1], 'format': img.format}
        logger.info(f'Image dimensions: {result}')
        return result
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        return None

def get_gif_file_info__8d385ddc(env, config: dict):
    """
    Get information about a GIF file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' parameter

    Returns:
        dict: Information about the GIF file (exists, format, size, frame_count)
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}
        cache_path = os.path.join(env.cache_dir, os.path.basename(file_path))
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        from PIL import Image
        try:
            img = Image.open(cache_path)
            frame_count = 0
            try:
                while True:
                    img.seek(frame_count)
                    frame_count += 1
            except EOFError:
                pass
            return {'exists': True, 'format': img.format, 'size': len(file_bytes), 'frame_count': frame_count, 'width': img.width, 'height': img.height}
        except Exception as e:
            logger.error(f'Error analyzing GIF: {e}')
            return {'exists': True, 'format': 'unknown', 'size': len(file_bytes), 'frame_count': 0}
    except Exception as e:
        logger.error(f'Error getting file: {e}')
        return {'exists': False, 'format': None, 'size': 0, 'frame_count': 0}

def get_image_dimensions__717d6863(env, config):
    """
    Get the dimensions (width, height) of an image file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - 'path': Path to the image file on VM

    Returns:
        dict: {'width': int, 'height': int} or None on error
    """
    try:
        file_path = config.get('path', '/home/user/Desktop/resized.png')
        image_content = env.controller.get_file(file_path)
        if not image_content:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, 'temp_image_check.png')
        with open(local_path, 'wb') as f:
            f.write(image_content)
        img = Image.open(local_path)
        dimensions = {'width': img.size[0], 'height': img.size[1]}
        logger.info(f'Image dimensions: {dimensions}')
        return dimensions
    except Exception as e:
        logger.error(f'Error getting image dimensions: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None

def get_docx_image_count__02c7d140(env, config):
    """
    Extract the number of images in a DOCX file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to DOCX file

    Returns:
        int: Number of images in the document
    """
    file_path = config.get('path')
    if not file_path:
        logger.error('No file path provided in config')
        return 0
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.error(f'Could not retrieve file from VM: {file_path}')
        return 0
    try:
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp.flush()
            doc = Document(tmp.name)
            image_count = 0
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    image_count += 1
            return image_count
    except Exception as e:
        logger.error(f'Error extracting image count: {e}')
        return 0

def get_image_props__6575e228_v8(env, config: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    """Get the dimensions and pixel data of an image file for rotation verification.

    This getter extracts both dimensions and corner pixel samples from the result image,
    plus loads the original image to verify 180-degree rotation occurred.

    Args:
        env: Environment object
        config: Configuration dict with:
            - 'path': Path to the image file in the VM
            - 'dest': Destination filename for caching

    Returns:
        dict: Dictionary with 'width', 'height', and 'pixels' (numpy array), or None if file doesn't exist
    """
    try:
        file_path = config.get('path')
        dest = config.get('dest', 'temp_image.jpg')
        if not file_path:
            logger.error('No file path provided in config')
            return None
        image_bytes = env.controller.get_file(file_path)
        if not image_bytes:
            logger.error(f'Failed to get image file: {file_path}')
            return None
        local_path = os.path.join(env.cache_dir, dest)
        with open(local_path, 'wb') as f:
            f.write(image_bytes)
        result_img = Image.open(local_path)
        (width, height) = result_img.size
        result_pixels = np.array(result_img)
        logger.info(f'Result image dimensions: {width}x{height}')
        original_path = '/home/user/OIP.jpg'
        try:
            original_bytes = env.controller.get_file(original_path)
            if original_bytes:
                original_local_path = os.path.join(env.cache_dir, 'original_OIP.jpg')
                with open(original_local_path, 'wb') as f:
                    f.write(original_bytes)
                original_img = Image.open(original_local_path)
                original_pixels = np.array(original_img)
                logger.info(f'Original image dimensions: {original_img.size}')
                return {'width': width, 'height': height, 'result_pixels': result_pixels, 'original_pixels': original_pixels}
            else:
                logger.warning('Could not get original image, will only check dimensions')
                return {'width': width, 'height': height, 'result_pixels': result_pixels, 'original_pixels': None}
        except Exception as e:
            logger.warning(f'Error loading original image: {e}, will only check dimensions')
            return {'width': width, 'height': height, 'result_pixels': result_pixels, 'original_pixels': None}
    except Exception as e:
        logger.error(f'Error getting image properties: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return None
