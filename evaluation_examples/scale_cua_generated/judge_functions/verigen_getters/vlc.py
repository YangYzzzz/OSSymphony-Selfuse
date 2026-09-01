"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_vlc_snapshot_file__57da68f2', 'get_vlc_bookmarks_status__18214e12ab9f5ae2ef34a7b82632e6f9', 'get_video_duration_file__da7e9a1f', 'get_audio_conversion_result__70ff6e4ac087c63cd4c29ec50188d44d', 'get_video_frame_at_timestamp__8d385ddc', 'get_vlc_screenshot_check__27395ab4949ce7b3fcdc03fed506b8b7', 'get_audio_file_info__4071abfac3d4e1d773938a9f2b9279b5', 'get_vlc_video_info__78c019ea', 'get_wallpaper_is_snapshot__b36f4827', 'get_snapshot_file_properties__4e252238', 'get_vlc_playback_rate__d5f08d8dc19ae133ba902b1660141393', 'get_vlc_playlist__aa13955b1d571b5f9f91c4b58d07755a', 'get_vlc_loop_status__869e58e4ccfbeafe3d06c94499c355f5', 'get_snapshot_file_exists__0706c584', 'get_snapshot_file_count__8588412f', 'get_vlc_config_and_file_status__8f080098', 'get_vlc_snapshot_exists__dd9409c8', 'get_vlc_volume_mute__95775285be787fb106c6c57a61517d2a', 'get_vlc_loop_status__9eea6b3721652abbf42884414d2268c7', 'get_ogg_audio_info__c03bf6861f4e39f709595e582520ac0f', 'get_vlc_max_volume_setting__dea0bbbbee03e4923af38de1331e256f', 'get_vlc_video_info__5e99ab87', 'get_audio_format_info__5d993657', 'get_vlc_mute_status__cf575cad41cf5daac3aae885ce8cfe2a', 'get_vlc_video_info__cb59ca2e', 'get_audio_duration__ee9c4304f47d297d44dedaad1e2983d6', 'get_vlc_volume__eef15ca1487d76cdd0b6405615ef653d', 'get_vlc_fullscreen_state__d7a7afc9', 'get_vlc_bookmark_accessibility__1c6b323b', 'get_vlc_playback_state__8f5f9872b32c7843d08f183a299b915a', 'get_snapshot_file_status__1ce093b98b48496fecfb155d31ed7704', 'get_vlc_bookmarks_count__3288b82e', 'get_desktop_snapshot_exists__d17b574f', 'get_audio_volume__28cc3b7e', 'get_vlc_continue_playback_config', 'get_mp3_exists__423224cbbe432d6315ffb9aa3c684c3a', 'get_vlc_playback_rate__a884bc65', 'get_vlc_screenshot_file__167faa79', 'get_vlc_video_info__fc9f20bb', 'get_video_rotation__ad793600129eb921cf29c68b343191af', 'get_vlc_video_info__bdd54294', 'get_vlc_is_playing_video__24795e62', 'get_vlc_video_info__2331e840', 'get_vlc_loop_mode__bd908fa8', 'get_wav_file_properties__580377a3f955e45d41d67d84cdd5fa88', 'get_vlc_repeat_loop_config', 'get_snapshot_format__a060a2ea', 'get_vlc_playback_time__2bf4fa278b8d46abde1b1e5c8949b96a', 'get_vlc_screenshot_path__6c9e8f1d4a3b2e5f7890abcd12345678', 'get_vlc_window_maximized__2c9788b1', 'get_vlc_volume__9cc567291e2f7ec7556a7b0f92f9b42b', 'get_vlc_volume__c48d7656', 'get_vlc_playback_rate__c979999ccfa7b4560c7eb8656e505ff9', 'get_snapshot_file_size__313dd5e1', 'get_vlc_video_info__636928c8', 'get_vlc_playback_state__df488c66748635e65016c0e9e7cc92c2', 'get_vlc_subtitle_track__65753b1a', 'get_vlc_bgcone_setting__ed0d0c08edef23629254f099c29e5e89', 'get_vlc_aspect_ratio__8a47ea01', 'get_vlc_aspect_ratio__d1470170', 'get_vlc_always_on_top__245c3f85', 'get_vlc_playback_speed__c3b737f4', 'get_audio_file_format__6475bf6e7aa0ef4599a6c11ec95f5406', 'get_vlc_video_effects__84f74170', 'get_vlc_playback_speed__b139c1b4', 'get_vlc_video_info__93a4e125', 'get_downloads_snapshot_exists__2cbf25da', 'get_vlc_video_info__a899663b', 'get_vlc_audio_muted__b5c6209b', 'get_mp3_file_info__ce2e74ee9437f284464037da6df4e453', 'get_desktop_mp4_files__04085b6d', 'get_vlc_snapshot_prefix__bdf11a45', 'get_video_file_exists__aec9e92c', 'get_mp3_count_in_dir__be089f93', 'get_vlc_always_on_top__b1199d7f2afb5cd10588e7e1e6cdc076', 'get_vlc_minimal_view_setting__99af05cec0829d789ac2d7bf7abe0481', 'get_vlc_video_info__6e7fd09f', 'get_vlc_playback_state__fc6f4242', 'get_video_file_properties__4939fb90', 'get_vlc_snapshot_directory__008263c0', 'get_audio_extraction_status__d0f84a157daee05fd0a57bb243a6ea5c', 'get_audio_file_properties__5d993657', 'get_vlc_loop_status__9e0433d6', 'get_video_metadata_file__bc3a25a7baab15992708f46f1d51e584', 'get_vlc_recordings_folder__055ce0bbe52cb194135b41b067eca986', 'get_vlc_playlist_count__fc73aaff', 'get_srt_and_video_status__ed96ceb6', 'get_mp3_audio_info__5e12822f6a68285e2467088dd7d25598']

def get_vlc_snapshot_file__57da68f2(env, config: Dict) -> Optional[str]:
    """
    Check if a VLC snapshot file exists in the Pictures directory.

    Note: The task config includes a preconfig step that cleans any existing
    vlcsnap-*.png files from ~/Pictures before task execution, ensuring that
    only newly created snapshots during task execution will be detected.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        The snapshot filename if found, None otherwise
    """
    pictures_path = '/home/user/Pictures'
    result = env.controller.run_bash_script(f"find {pictures_path} -maxdepth 1 -type f -name 'vlcsnap-*.png' 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output'].strip():
        logger.info('No VLC snapshot files found in Pictures directory')
        return None
    files = [f.strip() for f in result['output'].strip().split('\n') if f.strip()]
    snapshot_files = [os.path.basename(f) for f in files if os.path.basename(f).startswith('vlcsnap-') and f.endswith('.png')]
    if snapshot_files:
        logger.info(f'Found VLC snapshot file: {snapshot_files[0]}')
        return snapshot_files[0]
    logger.info('No VLC snapshot files found')
    return None

def get_vlc_bookmarks_status__18214e12ab9f5ae2ef34a7b82632e6f9(env, config: Dict[str, str]):
    """
    Checks if VLC bookmarks file was created.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'bookmarks_path' key

    Returns:
        dict: Bookmarks file status information
    """
    bookmarks_path = config.get('bookmarks_path', '/home/user/bookmarks.xspf')
    file_check = env.controller.run_bash_script(f"test -f {bookmarks_path} && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = 'EXISTS' in file_check.get('output', '')
    if not file_exists:
        return {'exists': False, 'content': '', 'has_video_entry': False}
    content_result = env.controller.run_bash_script(f'cat {bookmarks_path}', timeout=10)
    content = content_result.get('output', '')
    has_video_entry = 'video.mp4' in content
    return {'exists': True, 'content': content, 'has_video_entry': has_video_entry, 'path': bookmarks_path}

def get_video_duration_file__da7e9a1f(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the duration.txt file from the VM that should contain the video duration.

    Args:
        env: Environment object
        config: Configuration dict with 'duration_file' key specifying the path

    Returns:
        str: Path to the local cached file, or None if file doesn't exist
    """
    duration_file_path = config.get('duration_file', '/home/user/duration.txt')
    dest_filename = os.path.basename(duration_file_path)
    _path = os.path.join(env.cache_dir, dest_filename)
    try:
        file_content = env.controller.get_file(duration_file_path)
        if file_content is None:
            logger.warning(f'Duration file not found at: {duration_file_path}')
            return None
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(_path, 'wb') as f:
            f.write(file_content)
        logger.info(f'Successfully retrieved duration file: {_path} ({len(file_content)} bytes)')
        return _path
    except Exception as e:
        logger.error(f'Error retrieving duration file from {duration_file_path}: {e}')
        return None

def get_audio_conversion_result__70ff6e4ac087c63cd4c29ec50188d44d(env, config: Dict[str, str]):
    """
    Get information about converted audio file including format verification.
    Returns dict with file properties and format details.
    """
    file_path = config.get('path', '')
    command = f"\nimport os\nimport json\nimport subprocess\n\nfile_path = '{file_path}'\nresult = {{}}\n\nif os.path.exists(file_path):\n    result['exists'] = True\n    result['size'] = os.path.getsize(file_path)\n    result['extension'] = os.path.splitext(file_path)[1].lower()\n    result['basename'] = os.path.basename(file_path)\n\n    # Use file command to verify actual format\n    try:\n        file_output = subprocess.check_output(['file', '--mime-type', '-b', file_path], text=True).strip()\n        result['mime_type'] = file_output\n        result['is_audio'] = file_output.startswith('audio/')\n    except:\n        result['mime_type'] = 'unknown'\n        result['is_audio'] = False\nelse:\n    result['exists'] = False\n    result['size'] = 0\n    result['extension'] = ''\n    result['basename'] = ''\n    result['mime_type'] = ''\n    result['is_audio'] = False\n\nprint(json.dumps(result))\n"
    try:
        response = env.controller.execute_python_command(command)
        if response and response.get('output'):
            import json
            return json.loads(response['output'].strip())
        return {'exists': False, 'size': 0, 'extension': '', 'basename': '', 'mime_type': '', 'is_audio': False}
    except Exception as e:
        logger.error(f'Error checking audio conversion result: {e}')
        return {'exists': False, 'size': 0, 'extension': '', 'basename': '', 'mime_type': '', 'is_audio': False}

def get_video_frame_at_timestamp__8d385ddc(env, config):
    """
    Extract a frame from a video at a specific timestamp.

    Task: 2fe4b718-3bd7-46ec-bdce-b184f5653624
    Instruction: Extract a screenshot at exactly 3 seconds from the video 'src.mp4'
                 on the desktop and save it as 'screenshot_3s.jpg'.

    Args:
        env: Environment object
        config: Configuration dict with keys:
                - video_path: Path to the video file
                - timestamp: Timestamp in seconds to extract the frame

    Returns:
        str: Path to the extracted frame image file
    """
    video_path = config.get('video_path')
    timestamp = config.get('timestamp', 3.0)
    if not video_path:
        logger.error('video_path not specified in config')
        return None
    try:
        local_video_path = env.controller.get_file(video_path)
        if not local_video_path or not os.path.exists(local_video_path):
            logger.error(f'Video file not found: {video_path}')
            return None
        cap = cv2.VideoCapture(local_video_path)
        if not cap.isOpened():
            logger.error(f'Failed to open video file: {local_video_path}')
            return None
        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0
        logger.info(f'Video info - FPS: {fps}, Total frames: {total_frames}, Duration: {duration}s')
        if timestamp > duration:
            logger.warning(f'Requested timestamp {timestamp}s exceeds video duration {duration}s')
            cap.release()
            return None
        frame_number = int(timestamp * fps)
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_number)
        (ret, frame) = cap.read()
        cap.release()
        if not ret or frame is None:
            logger.error(f'Failed to read frame at timestamp {timestamp}s')
            return None
        frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
        image = Image.fromarray(frame_rgb)
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
        temp_path = temp_file.name
        temp_file.close()
        image.save(temp_path, 'JPEG')
        logger.info(f'Extracted frame at {timestamp}s from {video_path} -> {temp_path}')
        return temp_path
    except Exception as e:
        logger.error(f'Error extracting frame from video: {e}')
        return None

def get_vlc_screenshot_check__27395ab4949ce7b3fcdc03fed506b8b7(env, config: dict):
    """
    Check if a VLC screenshot exists with expected filename pattern.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File information including existence and basic properties
    """
    path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(path)
        if file_bytes is None:
            logger.warning(f'Screenshot file not found: {path}')
            return {'exists': False, 'filename': os.path.basename(path), 'size': 0}
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            img = Image.open(tmp_path)
            is_valid = img.format == 'PNG'
            (width, height) = img.size
            return {'exists': True, 'filename': os.path.basename(path), 'size': len(file_bytes), 'is_png': is_valid, 'width': width, 'height': height}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.error(f'Error checking screenshot {path}: {e}')
        return {'exists': False, 'filename': os.path.basename(path), 'size': 0, 'error': str(e)}

def get_audio_file_info__4071abfac3d4e1d773938a9f2b9279b5(env, config: Dict[str, str]):
    """
    Get information about an audio file on the VM.
    Returns dict with file existence, size, and extension.
    """
    file_path = config.get('path', '')
    command = f"\nimport os\nimport json\n\nfile_path = '{file_path}'\nresult = {{}}\n\nif os.path.exists(file_path):\n    result['exists'] = True\n    result['size'] = os.path.getsize(file_path)\n    result['extension'] = os.path.splitext(file_path)[1].lower()\n    result['filename'] = os.path.basename(file_path)\nelse:\n    result['exists'] = False\n    result['size'] = 0\n    result['extension'] = ''\n    result['filename'] = ''\n\nprint(json.dumps(result))\n"
    try:
        response = env.controller.execute_python_command(command)
        if response and response.get('output'):
            import json
            return json.loads(response['output'].strip())
        return {'exists': False, 'size': 0, 'extension': '', 'filename': ''}
    except Exception as e:
        logger.error(f'Error checking audio file: {e}')
        return {'exists': False, 'size': 0, 'extension': '', 'filename': ''}

def get_vlc_video_info__78c019ea(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.
    Also extracts frames from both videos for 180-degree rotation verification.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, ffprobe results, and frame hashes for comparison
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_Flipped.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    time.sleep(1)
    temp_dir = '/tmp/vlc_rotation_check'
    source_frame = f'{temp_dir}/source_frame.png'
    output_frame = f'{temp_dir}/output_frame.png'
    rotated_source_frame = f'{temp_dir}/source_frame_rotated.png'
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate,bit_rate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate,bit_rate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'", f'mkdir -p {temp_dir}', f"ffmpeg -y -i '{source_path}' -vf 'select=eq(n\\,100)' -vframes 1 '{source_frame}' 2>&1 | grep -E '(frame=|error|ERROR)' || echo 'FRAME_EXTRACT_DONE'", f"ffmpeg -y -i '{output_path}' -vf 'select=eq(n\\,100)' -vframes 1 '{output_frame}' 2>&1 | grep -E '(frame=|error|ERROR)' || echo 'FRAME_EXTRACT_DONE'", f"ffmpeg -y -i '{source_frame}' -vf 'transpose=2,transpose=2' '{rotated_source_frame}' 2>&1 | grep -E '(frame=|error|ERROR)' || echo 'ROTATION_DONE'", f"convert '{rotated_source_frame}' -resize 8x8! -colorspace Gray -format '%[fx:mean]' info: 2>/dev/null || echo 'ERROR'", f"convert '{output_frame}' -resize 8x8! -colorspace Gray -format '%[fx:mean]' info: 2>/dev/null || echo 'ERROR'", f"compare -metric RMSE '{rotated_source_frame}' '{output_frame}' null: 2>&1 || echo 'COMPARE_DONE'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4] + '---FRAME_EXTRACT---' + result_parts[5] + '---SOURCE_FRAME---' + result_parts[6] + '---OUTPUT_FRAME---' + result_parts[7] + '---ROTATION_DONE---' + result_parts[8] + '---ROTATED_HASH---' + result_parts[9] + '---OUTPUT_HASH---' + result_parts[10] + '---PIXEL_COMPARE---' + result_parts[11]
    return combined_result

def get_wallpaper_is_snapshot__b36f4827(env, config: dict):
    """Check if current wallpaper is a VLC snapshot file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        Boolean indicating if wallpaper is a snapshot
    """
    command = ['gsettings', 'get', 'org.gnome.desktop.background', 'picture-uri']
    result = env.controller.run_bash_script(' '.join(command), timeout=10)
    output = result.get('output', '').strip()
    if output.startswith("'file://"):
        output = output[8:-1]
    elif output.startswith('file://'):
        output = output[7:]
    elif output.startswith("'"):
        output = output[1:-1]
    if 'vlc-snap' not in output.lower():
        return False
    check_file_command = f"test -f '{output}' && echo 'exists' || echo 'not found'"
    file_check = env.controller.run_bash_script(check_file_command, timeout=10)
    if 'exists' not in file_check.get('output', ''):
        return False
    timestamp_command = f"stat -c '%Y' '{output}' 2>/dev/null || echo '0'"
    timestamp_result = env.controller.run_bash_script(timestamp_command, timeout=10)
    try:
        file_timestamp = int(timestamp_result.get('output', '0').strip())
        current_time = int(time.time())
        if file_timestamp > 0 and current_time - file_timestamp < 600:
            return True
    except (ValueError, TypeError):
        logger.warning('Could not verify file timestamp, using filename check only')
        return True
    return False

def get_snapshot_file_properties__4e252238(env, config: Dict[str, str]) -> Optional[Dict[str, any]]:
    """
    Gets properties of a snapshot file created from VLC.

    Args:
        env: Environment object
        config: Configuration dict containing:
            - snapshot_path: Path to the snapshot file on the VM

    Returns:
        Dict with file properties if file exists, None otherwise:
        {
            'exists': bool,
            'is_image': bool,
            'size': int (file size in bytes),
            'path': str (file path)
        }
    """
    snapshot_path = config.get('snapshot_path', '/home/user/snapshot.png')
    try:
        vm_ip = env.vm_ip
        port = env.server_port
        import requests
        check_cmd = ['python3', '-c', f"import os; path = '{snapshot_path}'; exists = os.path.exists(path); size = os.path.getsize(path) if exists else 0; print(f'{{exists}}|{{size}}')"]
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': check_cmd, 'shell': False})
        if response.status_code != 200:
            logger.error('Failed to check file properties. Status code: %d', response.status_code)
            return {'exists': False, 'is_image': False, 'size': 0, 'path': snapshot_path}
        output = response.json()['output'].strip()
        (exists_str, size_str) = output.split('|')
        exists = exists_str == 'True'
        size = int(size_str) if size_str.isdigit() else 0
        if not exists:
            return {'exists': False, 'is_image': False, 'size': 0, 'path': snapshot_path}
        check_image_cmd = ['python3', '-c', f"import os; path = '{snapshot_path}'; with open(path, 'rb') as f: magic = f.read(8); is_png = magic[:4] == b'\\x89PNG'; is_jpg = magic[:2] == b'\\xff\\xd8'; is_gif = magic[:3] == b'GIF'; is_bmp = magic[:2] == b'BM'; is_image = is_png or is_jpg or is_gif or is_bmp; print(is_image)"]
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': check_image_cmd, 'shell': False})
        if response.status_code != 200:
            logger.error('Failed to check image format. Status code: %d', response.status_code)
            is_image = False
        else:
            is_image_str = response.json()['output'].strip()
            is_image = is_image_str == 'True'
        return {'exists': exists, 'is_image': is_image, 'size': size, 'path': snapshot_path}
    except Exception as e:
        logger.error('Error getting snapshot file properties: %s', str(e))
        return {'exists': False, 'is_image': False, 'size': 0, 'path': snapshot_path}

def get_vlc_playback_rate__d5f08d8dc19ae133ba902b1660141393(env, config: Dict[str, str]):
    """
    Gets the current playback rate/speed from VLC's HTTP interface.
    Returns the path to the VLC status XML file.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlc_status.xml'))
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password))
        if response.status_code == 200:
            content = response.content
        else:
            logger.error('Failed to get vlc status. Status code: %d', response.status_code)
            return None
        with open(_path, 'wb') as f:
            f.write(content)
        return _path
    except Exception as e:
        logger.error(f'Error getting VLC playback rate: {e}')
        return None

def get_vlc_playlist__aa13955b1d571b5f9f91c4b58d07755a(env, config: dict):
    """Get VLC playlist file content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for playlist file

    Returns:
        dict: Playlist information including media items
    """
    file_path = config.get('path', '/home/user/.config/vlc/vlc-qt-interface.conf')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes:
            logger.warning(f'Could not read playlist file {file_path}')
            return {'media_items': []}
        content = file_bytes.decode('utf-8', errors='ignore')
        media_items = []
        for line in content.split('\n'):
            if 'recent-media' in line.lower() or 'playlist' in line.lower():
                media_items.append(line.strip())
        logger.info(f'Found {len(media_items)} playlist/recent entries')
        return {'media_items': media_items, 'raw_content': content}
    except Exception as e:
        logger.error(f'Error reading VLC playlist: {e}')
        return {'media_items': []}

def get_vlc_loop_status__869e58e4ccfbeafe3d06c94499c355f5(env, config: Dict[str, str]):
    """
    Gets the current loop/repeat status from VLC's HTTP interface.
    Returns the path to the VLC status XML file.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlc_status.xml'))
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password))
        if response.status_code == 200:
            content = response.content
        else:
            logger.error('Failed to get vlc status. Status code: %d', response.status_code)
            return None
        with open(_path, 'wb') as f:
            f.write(content)
        return _path
    except Exception as e:
        logger.error(f'Error getting VLC loop status: {e}')
        return None

def get_snapshot_file_exists__0706c584(env, config: dict):
    """Check if a snapshot/screenshot file exists at the specified path.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying the file to check

    Returns:
        bool: True if file exists, False otherwise
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f"test -f '{file_path}' && echo 'exists' || echo 'not_found'", timeout=10)
    output = result.get('output', '').strip()
    logger.info(f'Checking snapshot file at {file_path}: {output}')
    return output == 'exists'

def get_snapshot_file_count__8588412f(env, config: dict):
    """Count snapshot files in Pictures directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'directory' and 'pattern'

    Returns:
        Integer count of matching files
    """
    directory = config.get('directory', '/home/user/Pictures')
    pattern = config.get('pattern', 'vlc-snap*.png')
    command = f'ls {directory}/{pattern} 2>/dev/null | wc -l'
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        return 0

def get_vlc_config_and_file_status__8f080098(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get both VLC configuration and file existence status.

    Args:
        env: Environment object
        config: Configuration dict with 'vlc_dest' and 'file_path' keys

    Returns:
        dict: {
            'vlc_config': str (VLC config file content),
            'file_info': dict {'exists': bool, 'size': int or None}
        }
    """
    result = {'vlc_config': '', 'file_info': {'exists': False, 'size': None}}
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        logger.error(f'Unsupported operating system: {os_type}')
        return result
    try:
        _path = os.path.join(env.cache_dir, config.get('vlc_dest', 'vlcrc_result.txt'))
        content = env.controller.get_file(config_path)
        if content:
            with open(_path, 'wb') as f:
                f.write(content)
            result['vlc_config'] = content.decode('utf-8')
            logger.info(f'Successfully read VLC config from {config_path}')
        else:
            logger.warning(f'Could not read VLC config from {config_path}')
    except Exception as e:
        logger.error(f'Error reading VLC config: {e}')
    file_path = config.get('file_path', '')
    if file_path:
        try:
            vm_ip = env.vm_ip
            port = env.server_port
            command = f"[ -f '{file_path}' ] && stat -c '%s' '{file_path}' || echo 'NOT_EXISTS'"
            response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': True})
            if response.status_code == 200:
                output = response.json()['output'].strip()
                if output == 'NOT_EXISTS' or not output:
                    logger.info(f'File {file_path} does not exist')
                    result['file_info'] = {'exists': False, 'size': None}
                else:
                    try:
                        size = int(output)
                        logger.info(f'File {file_path} exists with size {size} bytes')
                        result['file_info'] = {'exists': True, 'size': size}
                    except ValueError:
                        logger.warning(f'Could not parse file size: {output}')
                        result['file_info'] = {'exists': False, 'size': None}
            else:
                logger.error(f'Failed to check file existence. Status code: {response.status_code}')
        except Exception as e:
            logger.error(f'Error checking file existence: {e}')
    return result

def get_vlc_snapshot_exists__dd9409c8(env, config: dict):
    """Check if a VLC snapshot file exists at the specified path.

    This getter searches for any vlc-snap*.png files in the Pictures directory,
    validates they are valid PNG images with content, and checks if they were
    created recently (within the last 5 minutes) to distinguish from pre-existing files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (used to extract directory)

    Returns:
        dict: {
            "exists": bool - whether a valid recent snapshot exists,
            "file_path": str - path to the found file (or None),
            "file_size": int - size in bytes (or 0),
            "is_recent": bool - whether file was created in last 5 minutes
        }
    """
    snapshot_dir = os.path.dirname(config.get('path', '/home/user/Pictures/vlc-snap.png'))
    find_command = f"find {snapshot_dir} -maxdepth 1 -name 'vlc-snap*.png' -type f"
    result = env.controller.run_bash_script(find_command, timeout=10)
    output = result.get('output', '').strip()
    if not output:
        logger.info('No VLC snapshot files found')
        return {'exists': False, 'file_path': None, 'file_size': 0, 'is_recent': False}
    files = [f.strip() for f in output.split('\n') if f.strip()]
    time_command = 'date +%s'
    time_result = env.controller.run_bash_script(time_command, timeout=10)
    try:
        current_time = int(time_result.get('output', '').strip())
    except (ValueError, AttributeError):
        logger.error('Failed to get current time from VM')
        current_time = 0
    time_threshold = current_time - 300
    valid_file = None
    valid_size = 0
    is_recent = False
    for file_path in files:
        stat_command = f"stat -c '%s %Y' '{file_path}' 2>/dev/null || echo '0 0'"
        stat_result = env.controller.run_bash_script(stat_command, timeout=10)
        stat_output = stat_result.get('output', '').strip()
        try:
            (size_str, mtime_str) = stat_output.split()
            file_size = int(size_str)
            file_mtime = int(mtime_str)
        except (ValueError, AttributeError):
            logger.warning(f'Failed to parse stat output for {file_path}')
            continue
        if file_size < 1024:
            logger.warning(f'File {file_path} is too small ({file_size} bytes)')
            continue
        png_check_command = f"file '{file_path}' | grep -q 'PNG image' && echo 'valid' || echo 'invalid'"
        png_result = env.controller.run_bash_script(png_check_command, timeout=10)
        png_output = png_result.get('output', '').strip()
        if png_output != 'valid':
            logger.warning(f'File {file_path} is not a valid PNG image')
            continue
        file_is_recent = file_mtime >= time_threshold
        if file_is_recent:
            valid_file = file_path
            valid_size = file_size
            is_recent = True
            logger.info(f'Found valid recent VLC snapshot: {file_path} ({file_size} bytes)')
            break
        if valid_file is None:
            valid_file = file_path
            valid_size = file_size
    return {'exists': valid_file is not None and is_recent, 'file_path': valid_file, 'file_size': valid_size, 'is_recent': is_recent}

def get_vlc_volume_mute__95775285be787fb106c6c57a61517d2a(env, config: Dict[str, str]):
    """
    Gets the current volume and mute status from VLC's HTTP interface.

    VLC's HTTP interface exposes volume and mute state via status.json/xml.
    This function handles version variability by checking multiple mute indicators:

    1. Query status.json for 'volume' and potential mute flags ('mute', 'muted', etc.)
    2. VLC reports volume as 0-256 (or 0-320) range
    3. When muted via UI button, VLC behavior:
       - Some versions: set explicit 'mute' field to true in status.json
       - All versions: set reported volume to 0 when mute button is pressed
    4. Fallback to XML parsing if JSON unavailable

    The function returns comprehensive state data to allow the metric to make
    an informed decision about mute status based on available indicators.

    Returns the path to a JSON file containing parsed VLC state information.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlc_status.xml'))
    json_path = _path.replace('.xml', '.json')
    try:
        url_json = f'http://{host}:{port}/requests/status.json'
        url_xml = f'http://{host}:{port}/requests/status.xml'
        result = {}
        try:
            response = requests.get(url_json, auth=('', password), timeout=5)
            if response.status_code == 200:
                status_data = response.json()
                result = {'volume': status_data.get('volume', 256), 'state': status_data.get('state', 'unknown'), 'source': 'json'}
                for mute_field in ['mute', 'muted', 'audio_mute', 'is_muted']:
                    if mute_field in status_data:
                        result['muted_explicit'] = bool(status_data[mute_field])
                        logger.info(f"Found explicit mute field '{mute_field}': {result['muted_explicit']}")
                        break
                logger.info(f"VLC status from JSON: volume={result['volume']}, state={result.get('state')}, muted_explicit={result.get('muted_explicit', 'N/A')}")
        except Exception as json_err:
            logger.warning(f'JSON endpoint failed, trying XML: {json_err}')
        if not result:
            response = requests.get(url_xml, auth=('', password), timeout=5)
            if response.status_code != 200:
                logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
                return None
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            volume_elem = tree.find('volume')
            state_elem = tree.find('state')
            result = {'volume': int(volume_elem.text) if volume_elem is not None and volume_elem.text else 256, 'state': state_elem.text if state_elem is not None and state_elem.text else 'unknown', 'source': 'xml'}
            for mute_attr in ['muted', 'mute', 'audio_mute']:
                if mute_attr in tree.attrib:
                    result['muted_explicit'] = tree.attrib[mute_attr].lower() in ('true', '1', 'yes')
                    logger.info(f"Found explicit mute attribute '{mute_attr}': {result['muted_explicit']}")
                    break
            if 'muted_explicit' not in result:
                for mute_elem_name in ['mute', 'muted', 'audio_mute']:
                    mute_elem = tree.find(mute_elem_name)
                    if mute_elem is not None and mute_elem.text:
                        result['muted_explicit'] = mute_elem.text.lower() in ('true', '1', 'yes')
                        logger.info(f"Found explicit mute element '{mute_elem_name}': {result['muted_explicit']}")
                        break
            logger.info(f"VLC status from XML: volume={result['volume']}, state={result.get('state')}, muted_explicit={result.get('muted_explicit', 'N/A')}")
        with open(json_path, 'w') as f:
            json.dump(result, f, indent=2)
        return json_path
    except requests.exceptions.RequestException as e:
        logger.error(f'Failed to connect to VLC HTTP interface: {e}')
        return None
    except Exception as e:
        logger.error(f'Error getting VLC volume/mute status: {e}')
        return None

def get_vlc_loop_status__9eea6b3721652abbf42884414d2268c7(env, config: Dict[str, str]):
    """
    Gets the current loop/repeat status from VLC's HTTP interface.
    Returns True if loop is enabled, False otherwise.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            loop_element = tree.find('loop')
            if loop_element is not None:
                loop_status = loop_element.text.lower() == 'true'
                logger.info(f'VLC loop status: {loop_status}')
                return loop_status
            else:
                logger.warning("Could not find 'loop' element in VLC status XML, assuming loop is off")
                return False
        else:
            logger.error('Failed to get VLC status. Status code: %d', response.status_code)
            return None
    except Exception as e:
        logger.error(f'Error getting VLC loop status: {e}')
        return None

def get_ogg_audio_info__c03bf6861f4e39f709595e582520ac0f(env, config: Dict[str, str]):
    """
    Get information about an OGG Vorbis audio file.
    Returns dict with existence, size, and format verification.
    """
    file_path = config.get('path', '')
    command = f"\nimport os\nimport json\nimport subprocess\n\nfile_path = '{file_path}'\nresult = {{}}\n\nif os.path.exists(file_path):\n    result['exists'] = True\n    result['size'] = os.path.getsize(file_path)\n    result['extension'] = os.path.splitext(file_path)[1].lower()\n    result['filename'] = os.path.basename(file_path)\n\n    # Check MIME type\n    try:\n        mime = subprocess.check_output(['file', '--mime-type', '-b', file_path], text=True).strip()\n        result['mime_type'] = mime\n        # OGG can be audio/ogg, audio/x-vorbis+ogg, application/ogg\n        result['is_ogg'] = 'ogg' in mime.lower()\n    except:\n        result['mime_type'] = 'unknown'\n        result['is_ogg'] = False\nelse:\n    result['exists'] = False\n    result['size'] = 0\n    result['extension'] = ''\n    result['filename'] = ''\n    result['mime_type'] = ''\n    result['is_ogg'] = False\n\nprint(json.dumps(result))\n"
    try:
        response = env.controller.execute_python_command(command)
        if response and response.get('output'):
            import json
            return json.loads(response['output'].strip())
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_ogg': False}
    except Exception as e:
        logger.error(f'Error checking OGG file: {e}')
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_ogg': False}

def get_vlc_max_volume_setting__dea0bbbbee03e4923af38de1331e256f(env, config: Dict[str, str]):
    """
    Gets the VLC maximum volume configuration setting from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' key for output file

    Returns:
        Path to the cached config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_video_info__5e99ab87(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, metadata, and corner pixel hashes
    to verify 180° rotation.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, ffprobe, and corner hash results
    """
    output_path = config.get('path', '/home/user/Desktop/Commercial_Upright.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffmpeg -ss 1 -i '{source_path}' -vframes 1 -vf 'crop=50:50:0:0' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{source_path}' -vframes 1 -vf 'crop=50:50:w-50:0' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{source_path}' -vframes 1 -vf 'crop=50:50:0:h-50' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{source_path}' -vframes 1 -vf 'crop=50:50:w-50:h-50' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{output_path}' -vframes 1 -vf 'crop=50:50:0:0' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{output_path}' -vframes 1 -vf 'crop=50:50:w-50:0' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{output_path}' -vframes 1 -vf 'crop=50:50:0:h-50' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'", f"ffmpeg -ss 1 -i '{output_path}' -vframes 1 -vf 'crop=50:50:w-50:h-50' -f framehash -hash md5 - 2>/dev/null | grep -E '^[0-9]' | head -1 || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4] + '---SRC_TL---' + result_parts[5] + '---SRC_TR---' + result_parts[6] + '---SRC_BL---' + result_parts[7] + '---SRC_BR---' + result_parts[8] + '---OUT_TL---' + result_parts[9] + '---OUT_TR---' + result_parts[10] + '---OUT_BL---' + result_parts[11] + '---OUT_BR---' + result_parts[12]
    return combined_result

def get_audio_format_info__5d993657(env, config: dict):
    """Get format information about an audio file using file command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Format information including file_type, is_mp3, file_exists, and filename_correct
    """
    file_path = config.get('path', '')
    expected_filename = 'Baby Justin Bieber.mp3'
    check_exists = env.controller.run_bash_script(f"test -f '{file_path}' && echo 'exists' || echo 'not_found'", timeout=10)
    file_exists = check_exists['returncode'] == 0 and check_exists['output'].strip() == 'exists'
    actual_filename = os.path.basename(file_path)
    filename_correct = actual_filename == expected_filename
    file_type = ''
    is_mp3 = False
    if file_exists:
        result = env.controller.run_bash_script(f"file -b '{file_path}'", timeout=10)
        if result['returncode'] == 0:
            file_type = result['output'].strip().lower()
            is_mp3 = 'audio' in file_type and ('mpeg' in file_type or 'mp3' in file_type or 'layer iii' in file_type)
    return {'file_type': file_type, 'is_mp3': is_mp3, 'file_exists': file_exists, 'filename_correct': filename_correct, 'path': file_path, 'expected_filename': expected_filename, 'actual_filename': actual_filename}

def get_vlc_mute_status__cf575cad41cf5daac3aae885ce8cfe2a(env, config: Dict[str, str]):
    """
    Gets the current mute status and playback state from VLC's HTTP interface.
    Returns a dict with 'is_muted' (bool) and 'state' (str), or None if unavailable.
    """
    import os
    import requests
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config['dest'])
    url = f'http://{host}:{port}/requests/status.xml'
    response = requests.get(url, auth=('', password))
    if response.status_code == 200:
        content = response.content
    else:
        logger.error('Failed to get vlc status. Status code: %d', response.status_code)
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    try:
        tree = ElementTree.fromstring(content)
        state_element = tree.find('state')
        playback_state = state_element.text if state_element is not None and state_element.text else 'stopped'
        information = tree.find('.//information')
        has_media = information is not None and len(information) > 0
        filename = None
        info_name = tree.find('.//info[@name="filename"]')
        if info_name is not None and info_name.text:
            filename = info_name.text
        is_muted = None
        mute_element = tree.find('mute')
        if mute_element is not None and mute_element.text:
            is_muted = mute_element.text.lower() in ['true', '1', 'yes']
        else:
            volume_element = tree.find('volume')
            if volume_element is not None and volume_element.text:
                volume = int(volume_element.text)
                is_muted = volume == 0
        result = {'is_muted': is_muted, 'state': playback_state, 'has_media': has_media, 'filename': filename}
        logger.info(f'VLC Status: {result}')
        return result
    except Exception as e:
        logger.error(f'Error parsing VLC status: {e}')
        return None

def get_vlc_video_info__cb59ca2e(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Ad_Rotated.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_audio_duration__ee9c4304f47d297d44dedaad1e2983d6(env, config: dict):
    """
    Get audio file duration using ffprobe (if available) or file size estimation.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        dict: {'exists': bool, 'duration': float, 'size': int}
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes or len(file_bytes) == 0:
        logger.info(f'File not found or empty: {file_path}')
        return {'exists': False, 'duration': 0.0, 'size': 0}
    file_size = len(file_bytes)
    escaped_path = file_path.replace("'", "'\\''")
    command = f"ffprobe -v error -show_entries format=duration -of default=noprint_wrappers=1:nokey=1 '{escaped_path}' 2>/dev/null || echo '0'"
    result = env.controller.run_bash_script(command, timeout=10)
    duration = 0.0
    if result and result.get('returncode') == 0:
        try:
            output = result.get('output', '').strip()
            duration = float(output)
            logger.info(f'Got duration from ffprobe: {duration}s')
        except (ValueError, TypeError):
            logger.warning(f"Could not parse ffprobe output: {result.get('output')}")
            duration = file_size / 16000.0
    logger.info(f'File: {file_path}, size: {file_size}, estimated duration: {duration}s')
    return {'exists': True, 'duration': duration, 'size': file_size}

def get_vlc_volume__eef15ca1487d76cdd0b6405615ef653d(env, config: Dict[str, str]) -> Optional[Dict[str, any]]:
    """
    Gets the current audio state from VLC's HTTP interface.
    Returns a dict with volume information and VLC process state.

    IMPORTANT LIMITATION:
    VLC's HTTP interface (status.xml) does NOT expose a separate mute boolean flag.
    When the user clicks the mute button in VLC, the player sets volume to 0 internally
    and remembers the previous volume for unmuting. However, the HTTP API only exposes
    the current volume level (0-256+), not the internal mute flag state.

    Therefore, this getter returns volume level, and the metric will infer mute state
    from volume == 0. This means:
    - If user clicks mute button: volume becomes 0 → detected as muted ✓
    - If user drags volume to 0: volume becomes 0 → detected as muted ✓
    Both actions are functionally equivalent from VLC's HTTP API perspective.

    This is a limitation of VLC's HTTP interface, not the evaluator implementation.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            volume_element = tree.find('volume')
            volume = None
            if volume_element is not None:
                try:
                    volume = int(volume_element.text)
                except (ValueError, TypeError):
                    logger.warning(f'Could not parse volume value: {volume_element.text}')
                    volume = None
            state_element = tree.find('state')
            playback_state = state_element.text if state_element is not None else None
            logger.debug(f"VLC status XML: {response.content.decode('utf-8')[:500]}")
            audio_state = {'volume': volume, 'playback_state': playback_state, 'vlc_responsive': True}
            logger.info(f'VLC audio state - volume: {volume}, playback_state: {playback_state}')
            return audio_state
        else:
            logger.error('Failed to get VLC status. Status code: %d', response.status_code)
            return {'vlc_responsive': False, 'volume': None, 'playback_state': None}
    except requests.exceptions.Timeout:
        logger.error('VLC HTTP interface timeout - VLC may not be running')
        return {'vlc_responsive': False, 'volume': None, 'playback_state': None}
    except requests.exceptions.ConnectionError:
        logger.error('Cannot connect to VLC HTTP interface - VLC may not be running')
        return {'vlc_responsive': False, 'volume': None, 'playback_state': None}
    except Exception as e:
        logger.error(f'Error getting VLC audio state: {e}')
        return {'vlc_responsive': False, 'volume': None, 'playback_state': None}

def get_vlc_fullscreen_state__d7a7afc9(env, config: dict):
    """
    Check if VLC is in fullscreen mode.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        bool: True if in fullscreen mode, False otherwise
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            root = ET.fromstring(response.content)
            fullscreen_element = root.find('.//fullscreen')
            if fullscreen_element is not None and fullscreen_element.text:
                return fullscreen_element.text.lower() in ['true', '1']
    except Exception as e:
        logger.error(f'Failed to get VLC fullscreen state: {e}')
    return False

def get_vlc_bookmark_accessibility__1c6b323b(env, config: dict):
    """
    Gets the VLC bookmark manager accessibility tree.
    The postconfig should have already opened VLC and triggered the bookmark dialog (Ctrl+B).

    Args:
        env: DesktopEnv instance
        config: Configuration dict (not used but required for signature)

    Returns:
        Dict containing accessibility tree information about the bookmark dialog
    """
    try:
        a11y_tree = env.controller.get_accessibility_tree()
        if not a11y_tree:
            logger.warning('No accessibility tree found')
            return {'has_dialog': False, 'bookmark_entries': [], 'error': 'No accessibility tree available'}
        bookmark_info = {'has_dialog': False, 'bookmark_entries': [], 'raw_tree': str(a11y_tree)[:1000]}
        tree_str = str(a11y_tree).lower()
        if 'bookmark' in tree_str or 'custom bookmark' in tree_str:
            bookmark_info['has_dialog'] = True
            import re
            time_patterns = re.findall('\\b\\d{1,2}:\\d{2}(?:\\.\\d+)?\\b|\\b\\d+\\.\\d+\\s*s\\b|\\b\\d+\\s*s\\b', tree_str)
            if time_patterns:
                bookmark_info['bookmark_entries'] = time_patterns
                logger.info(f'Found potential bookmark time entries: {time_patterns}')
        return bookmark_info
    except Exception as e:
        logger.error(f'Error getting VLC bookmark accessibility tree: {e}')
        return {'has_dialog': False, 'bookmark_entries': [], 'error': str(e)}

def get_vlc_playback_state__8f5f9872b32c7843d08f183a299b915a(env, config: Dict[str, str]):
    """
    Gets the current playback state from VLC's HTTP interface.
    Returns the state (playing, paused, stopped) from VLC status.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlc_status.xml'))
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password))
        if response.status_code == 200:
            content = response.content
        else:
            logger.error('Failed to get vlc status. Status code: %d', response.status_code)
            return None
        with open(_path, 'wb') as f:
            f.write(content)
        return _path
    except Exception as e:
        logger.error(f'Error getting VLC playback state: {e}')
        return None

def get_snapshot_file_status__1ce093b98b48496fecfb155d31ed7704(env, config: Dict[str, str]):
    """
    Checks if snapshot image was created from video.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'snapshot_path' key

    Returns:
        dict: Snapshot file status information
    """
    snapshot_path = config.get('snapshot_path', '/home/user/snapshot.png')
    file_check = env.controller.run_bash_script(f"test -f {snapshot_path} && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = 'EXISTS' in file_check.get('output', '')
    if not file_exists:
        return {'exists': False, 'size': 0, 'is_image': False}
    size_cmd = f'stat -c%s {snapshot_path}'
    size_result = env.controller.run_bash_script(size_cmd, timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
    except:
        size = 0
    file_cmd = f'file {snapshot_path}'
    file_result = env.controller.run_bash_script(file_cmd, timeout=10)
    file_output = file_result.get('output', '')
    is_image = 'image' in file_output.lower() or 'PNG' in file_output or 'JPEG' in file_output
    return {'exists': True, 'size': size, 'is_image': is_image, 'path': snapshot_path}

def get_vlc_bookmarks_count__3288b82e(env, config: dict):
    """
    Get the count of bookmarks in VLC by checking the bookmarks XSPF file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        dict: Contains 'count' and 'bookmarks' list with positions
    """
    bookmark_paths = ['/home/user/.local/share/vlc/ml.xspf', '/home/user/.config/vlc/bookmarks.xspf', '/home/user/.local/share/vlc/bookmarks.xspf']
    bookmarks = []
    for bookmark_path in bookmark_paths:
        result = env.controller.run_bash_script(f'cat {bookmark_path} 2>/dev/null', timeout=10)
        if result['returncode'] == 0 and result['output']:
            try:
                root = ET.fromstring(result['output'])
                ns = {'xspf': 'http://xspf.org/ns/0/'}
                tracks = root.findall('.//xspf:track', ns)
                if not tracks:
                    tracks = root.findall('.//track')
                for track in tracks:
                    bookmark_info = {'position': None, 'title': None}
                    title_elem = track.find('.//xspf:title', ns) or track.find('.//title')
                    if title_elem is not None:
                        bookmark_info['title'] = title_elem.text
                    ext_elem = track.find('.//xspf:extension', ns) or track.find('.//extension')
                    if ext_elem is not None:
                        for child in ext_elem:
                            if 'id' in child.tag or 'time' in child.tag.lower():
                                bookmark_info['position'] = child.text
                    location_elem = track.find('.//xspf:location', ns) or track.find('.//location')
                    if location_elem is not None and location_elem.text:
                        match = re.search('#t=([0-9.]+)', location_elem.text)
                        if match:
                            bookmark_info['position'] = match.group(1)
                    bookmarks.append(bookmark_info)
                if bookmarks:
                    break
            except ET.ParseError as e:
                logger.debug(f'Failed to parse {bookmark_path}: {e}')
                continue
            except Exception as e:
                logger.debug(f'Error processing {bookmark_path}: {e}')
                continue
    if not bookmarks:
        qt_conf_path = '/home/user/.config/vlc/vlc-qt-interface.conf'
        result = env.controller.run_bash_script(f"cat {qt_conf_path} 2>/dev/null | grep -E '^bookmarks=' | head -1", timeout=10)
        if result['returncode'] == 0 and result['output']:
            bookmark_line = result['output'].strip()
            if bookmark_line and '=' in bookmark_line:
                bookmark_values = bookmark_line.split('=', 1)[1]
                if bookmark_values:
                    positions = bookmark_values.split(',')
                    for pos in positions:
                        pos = pos.strip()
                        if pos:
                            bookmarks.append({'position': pos, 'title': None})
    if not bookmarks:
        vlcrc_path = '/home/user/.config/vlc/vlcrc'
        result = env.controller.run_bash_script(f"cat {vlcrc_path} 2>/dev/null | grep -E 'bookmark' | head -20", timeout=10)
        if result['returncode'] == 0 and result['output']:
            for line in result['output'].split('\n'):
                line = line.strip()
                match = re.search('bookmark.*=\\s*([0-9.]+)', line)
                if match:
                    bookmarks.append({'position': match.group(1), 'title': None})
    return {'count': len(bookmarks), 'bookmarks': bookmarks}

def get_desktop_snapshot_exists__d17b574f(env, config: dict):
    """Check if a snapshot file exists on Desktop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'filename' key

    Returns:
        Boolean indicating if file exists on Desktop
    """
    filename = config.get('filename', 'vlc-snap*.png')
    desktop_path = '/home/user/Desktop'
    command = f"ls {desktop_path}/{filename} 2>/dev/null && echo 'exists' || echo 'not_found'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    return 'exists' in output

def get_audio_volume__28cc3b7e(env, config):
    """
    Get the current audio output volume as an integer percentage.

    Handles multiple audio backends (PulseAudio, PipeWire, ALSA) and returns
    the volume of the default audio sink. Returns a clean integer value without
    newlines or formatting.

    Args:
        env: Environment object
        config: Configuration dict (unused, but required by signature)

    Returns:
        int: Volume percentage (0-100), or None if unable to determine
    """
    try:
        command = "pactl get-sink-volume @DEFAULT_SINK@ | grep -oP '\\d+%' | head -n 1 | sed 's/%//'"
        result = env.controller.execute(command)
        if result and result.strip():
            volume = int(result.strip())
            return volume
    except (ValueError, AttributeError):
        pass
    try:
        command = "pactl list sinks | grep -A 15 'State: RUNNING' | grep 'Volume:' | head -n 1 | awk '{print $5}' | sed 's/%//'"
        result = env.controller.execute(command)
        if result and result.strip():
            volume = int(result.strip())
            return volume
    except (ValueError, AttributeError):
        pass
    try:
        command = "amixer get Master | grep -oP '\\d+%' | head -n 1 | sed 's/%//'"
        result = env.controller.execute(command)
        if result and result.strip():
            volume = int(result.strip())
            return volume
    except (ValueError, AttributeError):
        pass
    return None

def get_vlc_continue_playback_config(env, config: Dict[str, str]):
    """
    Reads the VLC configuration file to check continue playback settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_mp3_exists__423224cbbe432d6315ffb9aa3c684c3a(env, config):
    """
    Check if MP3 file exists at specified path and validate it's a valid MP3 with audio content.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains 'exists', 'size', 'has_valid_format', and 'duration' information
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f"if [ -f '{file_path}' ]; then echo 'true'; else echo 'false'; fi", timeout=10)
    exists = result.get('output', '').strip() == 'true'
    if not exists:
        logger.info(f'File {file_path} does not exist')
        return {'exists': False, 'size': 0, 'has_valid_format': False, 'duration': 0}
    size_result = env.controller.run_bash_script(f"stat -c %s '{file_path}' 2>/dev/null || stat -f %z '{file_path}' 2>/dev/null", timeout=10)
    file_size = 0
    try:
        file_size = int(size_result.get('output', '0').strip())
    except ValueError:
        logger.warning(f'Could not determine file size for {file_path}')
    magic_check = env.controller.run_bash_script(f"head -c 3 '{file_path}' | xxd -p", timeout=10)
    magic_bytes = magic_check.get('output', '').strip()
    has_valid_format = magic_bytes.startswith('494433') or magic_bytes.startswith('fffb') or magic_bytes.startswith('fff3') or magic_bytes.startswith('fff2')
    duration = 0
    duration_result = env.controller.run_bash_script(f"ffprobe -v error -show_entries format=duration -of default=noprint_wrappers=1:nokey=1 '{file_path}' 2>/dev/null", timeout=15)
    try:
        duration_str = duration_result.get('output', '0').strip()
        if duration_str and duration_str != 'N/A':
            duration = float(duration_str)
    except (ValueError, AttributeError):
        logger.warning(f'Could not determine audio duration for {file_path}')
    logger.info(f'File {file_path} - exists: {exists}, size: {file_size}, valid_format: {has_valid_format}, duration: {duration}')
    return {'exists': exists, 'size': file_size, 'has_valid_format': has_valid_format, 'duration': duration}

def get_vlc_playback_rate__a884bc65(env, config: Dict[str, str]):
    """
    Gets the current playback rate/speed from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        float: Current playback rate (1.0 = normal speed, 0.5 = half speed, 2.0 = double speed)
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            rate_element = tree.find('rate')
            if rate_element is not None and rate_element.text:
                rate = float(rate_element.text)
                logger.info(f'VLC Playback Rate: {rate}')
                return rate
            else:
                logger.warning('Rate element not found in VLC status XML')
                return None
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC playback rate: {e}')
        return None

def get_vlc_screenshot_file__167faa79(env, config):
    """
    Check if VLC screenshot file exists on Desktop using command execution.

    Args:
        env: DesktopEnv instance
        config: Dict with 'pattern' for filename pattern to search

    Returns:
        str: Filename found, or empty string if not found
    """
    pattern = config.get('pattern', 'vlc-snap*.png')
    try:
        vm_ip = env.vm_ip
        port = env.server_port
        command = ['bash', '-c', f'ls -1 /home/user/Desktop/{pattern} 2>/dev/null | head -1 | xargs -r basename']
        response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': command, 'shell': False})
        if response.status_code == 200:
            result = response.json()['output'].strip()
            logger.info(f"VLC screenshot search result: '{result}'")
            return result
        else:
            logger.error(f'Failed to execute command. Status code: {response.status_code}')
            return ''
    except Exception as e:
        logger.error(f'Error checking screenshot: {e}')
        return ''

def get_vlc_video_info__fc9f20bb(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Macintosh_Commercial_Fixed.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_video_rotation__ad793600129eb921cf29c68b343191af(env, config: dict):
    """Get video rotation metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key specifying video file path on VM

    Returns:
        dict: {
            'exists': bool,
            'path': str,
            'rotation': int (0, 90, 180, 270),
            'width': int,
            'height': int,
            'file_size': int
        }
    """
    file_path = config.get('path', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if not file_bytes or len(file_bytes) == 0:
            logger.warning(f'Video file {file_path} does not exist or is empty')
            return {'exists': False, 'path': file_path, 'rotation': None, 'width': None, 'height': None, 'file_size': 0}
        with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            cmd = ['ffprobe', '-v', 'quiet', '-print_format', 'json', '-show_streams', '-show_format', tmp_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            if result.returncode != 0:
                logger.error(f'ffprobe failed for {file_path}: {result.stderr}')
                return {'exists': True, 'path': file_path, 'rotation': 0, 'width': None, 'height': None, 'file_size': len(file_bytes)}
            metadata = json.loads(result.stdout)
            video_stream = None
            for stream in metadata.get('streams', []):
                if stream.get('codec_type') == 'video':
                    video_stream = stream
                    break
            if not video_stream:
                logger.warning(f'No video stream found in {file_path}')
                return {'exists': True, 'path': file_path, 'rotation': 0, 'width': None, 'height': None, 'file_size': len(file_bytes)}
            rotation = 0
            for side_data in video_stream.get('side_data_list', []):
                if side_data.get('side_data_type') == 'Display Matrix':
                    rotation = side_data.get('rotation', 0)
                    break
            if rotation == 0:
                tags = video_stream.get('tags', {})
                rotation = int(tags.get('rotate', 0))
            rotation = int(rotation) % 360
            width = video_stream.get('width')
            height = video_stream.get('height')
            logger.info(f'Video {file_path}: rotation={rotation}, width={width}, height={height}, size={len(file_bytes)}')
            return {'exists': True, 'path': file_path, 'rotation': rotation, 'width': width, 'height': height, 'file_size': len(file_bytes)}
        finally:
            try:
                os.unlink(tmp_path)
            except Exception as e:
                logger.warning(f'Failed to delete temp file {tmp_path}: {e}')
    except subprocess.TimeoutExpired:
        logger.error(f'ffprobe timed out for {file_path}')
        return {'exists': True, 'path': file_path, 'rotation': None, 'width': None, 'height': None, 'file_size': 0}
    except Exception as e:
        logger.error(f'Error getting video metadata for {file_path}: {e}')
        return {'exists': False, 'path': file_path, 'rotation': None, 'width': None, 'height': None, 'file_size': 0}

def get_vlc_video_info__bdd54294(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_270.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_vlc_is_playing_video__24795e62(env, config):
    """Get VLC playback status and check if playing a specific video file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' for output file

    Returns:
        str: Path to VLC status XML file
    """
    import requests
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config['dest'])
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            content = response.content
        else:
            logger.error('Failed to get vlc status. Status code: %d', response.status_code)
            return None
    except Exception as e:
        logger.error(f'Failed to connect to VLC: {e}')
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_video_info__2331e840(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_Rotated_90.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'", f"ls -la '{source_path}'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4] + '---SOURCE_LS---' + (result_parts[5] if len(result_parts) > 5 else 'ERROR')
    return combined_result

def get_vlc_loop_mode__bd908fa8(env, config: dict):
    """
    Check if VLC loop/repeat mode is enabled by reading config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        dict: Dictionary with loop mode settings
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    result = env.controller.run_bash_script(f"cat {config_path} 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning('Could not read VLC config file')
        return {'loop_enabled': False, 'repeat_enabled': False}
    config_content = result['output']
    settings = {'loop_enabled': False, 'repeat_enabled': False}
    for line in config_content.split('\n'):
        line = line.strip()
        if line.startswith('#') or not line:
            continue
        if '=' in line:
            (key, value) = line.split('=', 1)
            key = key.strip()
            value = value.strip()
            if 'loop' in key.lower() and value == '1':
                settings['loop_enabled'] = True
            if 'repeat' in key.lower() and value == '1':
                settings['repeat_enabled'] = True
    return settings

def get_wav_file_properties__580377a3f955e45d41d67d84cdd5fa88(env, config: Dict[str, str]):
    """
    Get properties of a WAV audio file.
    Returns dict with file info and WAV-specific properties.
    """
    file_path = config.get('path', '')
    command = f"\nimport os\nimport json\nimport subprocess\n\nfile_path = '{file_path}'\nresult = {{}}\n\nif os.path.exists(file_path):\n    result['exists'] = True\n    result['size'] = os.path.getsize(file_path)\n    result['extension'] = os.path.splitext(file_path)[1].lower()\n    result['filename'] = os.path.basename(file_path)\n\n    # Verify MIME type\n    try:\n        mime = subprocess.check_output(['file', '--mime-type', '-b', file_path], text=True).strip()\n        result['mime_type'] = mime\n        result['is_wav'] = mime in ['audio/x-wav', 'audio/wav']\n    except:\n        result['mime_type'] = 'unknown'\n        result['is_wav'] = False\nelse:\n    result['exists'] = False\n    result['size'] = 0\n    result['extension'] = ''\n    result['filename'] = ''\n    result['mime_type'] = ''\n    result['is_wav'] = False\n\nprint(json.dumps(result))\n"
    try:
        response = env.controller.execute_python_command(command)
        if response and response.get('output'):
            import json
            return json.loads(response['output'].strip())
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_wav': False}
    except Exception as e:
        logger.error(f'Error checking WAV file: {e}')
        return {'exists': False, 'size': 0, 'extension': '', 'filename': '', 'mime_type': '', 'is_wav': False}

def get_vlc_repeat_loop_config(env, config: Dict[str, str]):
    """
    Reads the VLC configuration file to check repeat/loop settings.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_snapshot_format__a060a2ea(env, config: dict):
    """Get snapshot file format/type.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path'

    Returns:
        String with file format (e.g., 'PNG', 'JPEG')
    """
    filepath = config.get('path', '/home/user/Pictures/vlc-snap*.png')
    if '*' in filepath:
        command = f"file -b $(ls {filepath} 2>/dev/null | head -1) 2>/dev/null || echo 'not_found'"
    else:
        command = f"file -b {filepath} 2>/dev/null || echo 'not_found'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'PNG' in output:
        return 'PNG'
    elif 'JPEG' in output or 'JPG' in output:
        return 'JPEG'
    else:
        return 'UNKNOWN'

def get_vlc_playback_time__2bf4fa278b8d46abde1b1e5c8949b96a(env, config: Dict[str, str]):
    """
    Gets the current playback time position from VLC's HTTP interface.
    Returns the path to the VLC status XML file.
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config.get('dest', 'vlc_status.xml'))
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password))
        if response.status_code == 200:
            content = response.content
        else:
            logger.error('Failed to get vlc status. Status code: %d', response.status_code)
            return None
        with open(_path, 'wb') as f:
            f.write(content)
        return _path
    except Exception as e:
        logger.error(f'Error getting VLC playback time: {e}')
        return None

def get_vlc_screenshot_path__6c9e8f1d4a3b2e5f7890abcd12345678(env, config: Dict[str, str]):
    """
    Gets the VLC screenshot path configuration.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' for output filename

    Returns:
        Path to the saved config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_window_maximized__2c9788b1(env, config: Dict[str, str]):
    """
    Check if VLC window is maximized (but not fullscreen).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with app_class_name

    Returns:
        dict: Dictionary with 'is_maximized' (bool) and 'is_fullscreen' (bool)
    """
    try:
        window_size = env.controller.get_vm_window_size(app_class_name=config.get('app_class_name', 'vlc'))
        screen_size = env.controller.get_vm_screen_size()
        if window_size is None or screen_size is None:
            logger.warning('Could not get window or screen size')
            return {'is_maximized': False, 'is_fullscreen': False}
        is_fullscreen = window_size.get('width') == screen_size.get('width') and window_size.get('height') == screen_size.get('height')
        width_margin = 50
        height_margin = 100
        is_maximized = abs(window_size.get('width', 0) - screen_size.get('width', 0)) <= width_margin and abs(window_size.get('height', 0) - screen_size.get('height', 0)) <= height_margin
        result = {'is_maximized': is_maximized and (not is_fullscreen), 'is_fullscreen': is_fullscreen, 'window_width': window_size.get('width', 0), 'window_height': window_size.get('height', 0), 'screen_width': screen_size.get('width', 0), 'screen_height': screen_size.get('height', 0)}
        logger.info(f'VLC Window Status: {result}')
        return result
    except Exception as e:
        logger.error(f'Error checking VLC window maximized status: {e}')
        return {'is_maximized': False, 'is_fullscreen': False}

def get_vlc_volume__9cc567291e2f7ec7556a7b0f92f9b42b(env, config: Dict[str, str]):
    """
    Gets the current VLC status including volume, playback state, and current media.
    Returns a dict with volume, state, and filename information.
    """
    import os
    import requests
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    _path = os.path.join(env.cache_dir, config['dest'])
    url = f'http://{host}:{port}/requests/status.xml'
    response = requests.get(url, auth=('', password))
    if response.status_code == 200:
        content = response.content
    else:
        logger.error('Failed to get vlc status. Status code: %d', response.status_code)
        return None
    with open(_path, 'wb') as f:
        f.write(content)
    try:
        tree = ElementTree.fromstring(content)
        volume_element = tree.find('volume')
        volume = int(volume_element.text) if volume_element is not None and volume_element.text else None
        state_element = tree.find('state')
        state = state_element.text if state_element is not None and state_element.text else None
        filename = None
        information = tree.find('information')
        if information is not None:
            for category in information.findall('category'):
                if category.get('name') == 'meta':
                    for info in category.findall('info'):
                        if info.get('name') == 'filename':
                            filename = info.text
                            break
                    break
        result = {'volume': volume, 'state': state, 'filename': filename}
        logger.info(f'VLC Status - Volume: {volume}, State: {state}, Filename: {filename}')
        return result
    except Exception as e:
        logger.error(f'Error parsing VLC status: {e}')
        return None

def get_vlc_volume__c48d7656(env, config: Dict[str, str]):
    """
    Gets the current volume level from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' key for output path

    Returns:
        int: Current volume level (0-512, where 256 is 100%)
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            volume_element = tree.find('volume')
            if volume_element is not None and volume_element.text:
                volume = int(volume_element.text)
                logger.info(f'VLC Volume: {volume}')
                return volume
            else:
                logger.warning('Volume element not found in VLC status XML')
                return None
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC volume: {e}')
        return None

def get_vlc_playback_rate__c979999ccfa7b4560c7eb8656e505ff9(env, config: Dict[str, str]):
    """
    Gets the current playback rate/speed from VLC's HTTP interface.
    Returns the playback rate as a float (1.0 = normal speed).
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            rate_element = tree.find('rate')
            if rate_element is not None:
                rate = float(rate_element.text)
                logger.info(f'VLC playback rate: {rate}')
                return rate
            else:
                logger.warning("Could not find 'rate' element in VLC status XML, assuming normal speed")
                return 1.0
        else:
            logger.error('Failed to get VLC status. Status code: %d', response.status_code)
            return None
    except Exception as e:
        logger.error(f'Error getting VLC playback rate: {e}')
        return None

def get_snapshot_file_size__313dd5e1(env, config: dict):
    """Get size of snapshot file in bytes.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path'

    Returns:
        Integer file size in bytes, or 0 if file doesn't exist
    """
    filepath = config.get('path', '/home/user/Pictures/vlc-snap*.png')
    if '*' in filepath:
        command = f"ls -l {filepath} 2>/dev/null | head -1 | awk '{{print $5}}'"
    else:
        command = f"stat -c %s {filepath} 2>/dev/null || echo '0'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    try:
        return int(output)
    except ValueError:
        return 0

def get_vlc_video_info__636928c8(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_Rotated_CCW.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_vlc_playback_state__df488c66748635e65016c0e9e7cc92c2(env, config: Dict[str, str]):
    """
    Gets the current playback state from VLC's HTTP interface.
    Returns the state as a string (e.g., 'playing', 'paused', 'stopped').
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            state_element = tree.find('state')
            if state_element is not None:
                state = state_element.text
                logger.info(f'VLC playback state: {state}')
                return state
            else:
                logger.error("Could not find 'state' element in VLC status XML")
                return None
        else:
            logger.error('Failed to get VLC status. Status code: %d', response.status_code)
            return None
    except Exception as e:
        logger.error(f'Error getting VLC playback state: {e}')
        return None

def get_vlc_subtitle_track__65753b1a(env, config: Dict[str, str]):
    """
    Gets the current subtitle track ID from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        int: Current subtitle track ID (-1 or -2 means disabled, >= 0 means enabled), None on error
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            subtitle_element = tree.find('.//category[@name="meta"]/info[@name="subtitle"]')
            if subtitle_element is None:
                subtitle_element = tree.find('subtitles')
            if subtitle_element is not None and subtitle_element.text:
                subtitle_track = int(subtitle_element.text)
                logger.info(f'VLC Subtitle Track: {subtitle_track}')
                return subtitle_track
            else:
                logger.info('Subtitle element not found, assuming disabled (-1)')
                return -1
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC subtitle track: {e}')
        return None

def get_vlc_bgcone_setting__ed0d0c08edef23629254f099c29e5e89(env, config: Dict[str, str]):
    """
    Gets the VLC background cone configuration setting from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' key for output file

    Returns:
        Path to the cached config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_aspect_ratio__8a47ea01(env, config: dict):
    """
    Get VLC's current aspect ratio setting from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        str: Current aspect ratio setting or 'default'
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    result = env.controller.run_bash_script(f"cat {config_path} 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning('Could not read VLC config file')
        return 'default'
    config_content = result['output']
    aspect_ratio = None
    for line in config_content.split('\n'):
        line = line.strip()
        if line.startswith('#') or not line:
            continue
        if '=' in line:
            (key, value) = line.split('=', 1)
            key = key.strip()
            value = value.strip()
            if key in ['aspect-ratio', 'qt-video-aspect-ratio', 'custom-aspect-ratio']:
                if value and value not in ['', '0', '0.000000']:
                    aspect_ratio = value
                    break
            elif 'aspect-ratio' in key.lower():
                if value and value not in ['', '0', '0.000000']:
                    aspect_ratio = value
    if aspect_ratio:
        aspect_ratio = aspect_ratio.strip('"').strip("'")
        return aspect_ratio
    return 'default'

def get_vlc_aspect_ratio__d1470170(env, config: Dict[str, str]):
    """
    Gets the current aspect ratio setting from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: Current aspect ratio (e.g., 'default', '16:9', '4:3', '16:10', '1:1', '5:4')
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            from xml.etree import ElementTree
            tree = ElementTree.fromstring(response.content)
            aspect_paths = ['information/category[@name="Video"]/info[@name="Display aspect ratio"]', 'information/category[@name="Stream 0"]/info[@name="Display aspect ratio"]', 'aspectratio']
            for path in aspect_paths:
                element = tree.find(path)
                if element is not None and element.text:
                    aspect_ratio = element.text.strip()
                    logger.info(f'VLC Aspect Ratio: {aspect_ratio}')
                    return aspect_ratio
            logger.info("Aspect ratio not found in status XML, returning 'default'")
            return 'default'
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC aspect ratio: {e}')
        return None

def get_vlc_always_on_top__245c3f85(env, config: dict):
    """
    Check if VLC's always-on-top setting is enabled.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        bool: True if always-on-top is enabled, False otherwise
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    result = env.controller.run_bash_script(f"cat {config_path} 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning('Could not read VLC config file')
        return False
    config_content = result['output']
    for line in config_content.split('\n'):
        line = line.strip()
        if line.startswith('#') or not line:
            continue
        if '=' in line:
            (key, value) = line.split('=', 1)
            key = key.strip()
            value = value.strip()
            if 'on-top' in key.lower() or 'video-on-top' in key.lower():
                if value == '1':
                    return True
    return False

def get_vlc_playback_speed__c3b737f4(env, config: dict):
    """
    Get current VLC playback speed from HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        float: Current playback speed (1.0 = normal, 0.5 = half speed, 2.0 = double speed)
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            root = ET.fromstring(response.content)
            rate_element = root.find('.//rate')
            if rate_element is not None and rate_element.text:
                return float(rate_element.text)
    except Exception as e:
        logger.error(f'Failed to get VLC playback speed: {e}')
    return 1.0

def get_audio_file_format__6475bf6e7aa0ef4599a6c11ec95f5406(env, config: dict):
    """
    Check if an audio file exists and return its format information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key (VM file path)

    Returns:
        dict: {'exists': bool, 'format': str, 'size': int} or None if error
    """
    file_path = config.get('path', '')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        logger.info(f'File not found: {file_path}')
        return {'exists': False, 'format': None, 'size': 0}
    file_size = len(file_bytes)
    format_type = None
    if file_bytes[:3] == b'ID3' or (len(file_bytes) > 1 and file_bytes[0] == 255 and (file_bytes[1] & 224 == 224)):
        format_type = 'mp3'
    elif file_bytes[:4] == b'RIFF' and file_bytes[8:12] == b'WAVE':
        format_type = 'wav'
    elif file_bytes[:4] == b'OggS':
        format_type = 'ogg'
    elif file_bytes[:4] == b'fLaC':
        format_type = 'flac'
    logger.info(f'File found: {file_path}, size: {file_size}, format: {format_type}')
    return {'exists': True, 'format': format_type, 'size': file_size}

def get_vlc_video_effects__84f74170(env, config: dict):
    """
    Read VLC configuration file and extract mirror/horizontal flip video effect settings.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        dict: Dictionary with mirror effect status, including:
            - mirror_enabled: bool, whether mirror/horizontal flip is enabled
            - video_filter: str, the video filter setting
            - transform_type: str, the transform type if set
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    result = env.controller.run_bash_script(f"cat {config_path} 2>/dev/null || echo ''", timeout=10)
    if result['returncode'] != 0 or not result['output']:
        logger.warning('Could not read VLC config file')
        return {'mirror_enabled': False}
    config_content = result['output']
    settings = {'mirror_enabled': False, 'video_filter': None, 'transform_type': None}
    has_transform_filter = False
    has_mirror_filter = False
    transform_type_is_hflip = False
    has_mirror_mode_horizontal = False
    for line in config_content.split('\n'):
        line = line.strip()
        if line.startswith('#') or not line:
            continue
        if '=' in line:
            (key, value) = line.split('=', 1)
            key = key.strip()
            value = value.strip()
            if key == 'video-filter' or key == 'vout-filter':
                settings['video_filter'] = value
                if 'transform' in value.lower():
                    has_transform_filter = True
                if 'mirror' in value.lower():
                    has_mirror_filter = True
            elif key == 'transform-type':
                settings['transform_type'] = value
                if value in ['hflip', '1']:
                    transform_type_is_hflip = True
            elif key == 'mirror-mode' or key == 'mirror-direction':
                if value in ['horizontal', 'h', '1']:
                    has_mirror_mode_horizontal = True
    if has_transform_filter and transform_type_is_hflip or has_mirror_filter or has_mirror_mode_horizontal:
        settings['mirror_enabled'] = True
    logger.info(f"VLC mirror detection: transform_filter={has_transform_filter}, transform_type_hflip={transform_type_is_hflip}, mirror_filter={has_mirror_filter}, mirror_mode_horizontal={has_mirror_mode_horizontal}, result={settings['mirror_enabled']}")
    return settings

def get_vlc_playback_speed__b139c1b4(env, config: Dict[str, str]):
    """
    Gets the playback speed from VLC's configuration file.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        float: Playback speed (e.g., 1.0 for normal, 1.5 for 1.5x speed)
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    content = env.controller.get_file(config_path)
    config_content = content.decode('utf-8')
    playback_speed = 1.0
    for line in config_content.split('\n'):
        if line.startswith('#') or not line.strip():
            continue
        if line.startswith('rate='):
            try:
                rate_str = line.split('=', 1)[1].strip()
                playback_speed = float(rate_str)
                logger.info(f'Found VLC playback rate: {playback_speed}')
                break
            except (ValueError, IndexError) as e:
                logger.warning(f'Failed to parse playback rate from line: {line}, error: {e}')
    return playback_speed

def get_vlc_video_info__93a4e125(env, config: Dict[str, str]):
    """
    Get video file information and verify horizontal flip by extracting and comparing
    pixel data from mirrored positions.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, ffprobe, and pixel comparison results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_Mirrored.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,codec_name,duration,bit_rate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,codec_name,duration,bit_rate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffmpeg -i '{source_path}' -vf 'select=eq(n\\,0),crop=10:ih:0:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{output_path}' -vf 'select=eq(n\\,0),crop=10:ih:iw-10:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{source_path}' -vf 'select=eq(n\\,0),crop=10:ih:iw-10:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{output_path}' -vf 'select=eq(n\\,0),crop=10:ih:0:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{source_path}' -vf 'select=eq(n\\,30),crop=10:ih:0:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{output_path}' -vf 'select=eq(n\\,30),crop=10:ih:iw-10:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{source_path}' -vf 'select=eq(n\\,30),crop=10:ih:iw-10:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'", f"ffmpeg -i '{output_path}' -vf 'select=eq(n\\,30),crop=10:ih:0:0' -vframes 1 -f rawvideo -pix_fmt rgb24 - 2>/dev/null | md5sum || echo 'HASH_ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4] + '---EDGE_DATA---' + result_parts[5] + '---SRC_LEFT_F0---' + result_parts[6] + '---OUT_RIGHT_F0---' + result_parts[7] + '---SRC_RIGHT_F0---' + result_parts[8] + '---OUT_LEFT_F0---' + result_parts[9] + '---SRC_LEFT_F30---' + result_parts[10] + '---OUT_RIGHT_F30---' + result_parts[11] + '---SRC_RIGHT_F30---' + result_parts[12] + '---OUT_LEFT_F30---'
    return combined_result

def get_downloads_snapshot_exists__2cbf25da(env, config: dict):
    """Check if snapshot file exists in Downloads directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'pattern'

    Returns:
        Boolean indicating if file exists in Downloads
    """
    pattern = config.get('pattern', '*.png')
    downloads_path = '/home/user/Downloads'
    command = f"ls {downloads_path}/{pattern} 2>/dev/null && echo 'exists' || echo 'not_found'"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    return 'exists' in output

def get_vlc_video_info__a899663b(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/Apple_Commercial_Rotated_90.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_vlc_audio_muted__b5c6209b(env, config: Dict[str, str]):
    """
    Gets whether audio is muted in VLC.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        bool: True if audio is muted, False if not muted, None on error
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            volume_element = tree.find('volume')
            if volume_element is not None and volume_element.text:
                volume = int(volume_element.text)
                is_muted = volume == 0
                logger.info(f'VLC Audio Muted: {is_muted} (volume={volume})')
                return is_muted
            else:
                logger.warning('Volume element not found in VLC status XML')
                return None
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC audio mute status: {e}')
        return None

def get_mp3_file_info__ce2e74ee9437f284464037da6df4e453(env, config):
    """
    Get MP3 file information including existence, type, and basic properties.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File information including exists, is_audio, size
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f"if [ -f '{file_path}' ]; then echo 'EXISTS'; else echo 'NOT_EXISTS'; fi", timeout=10)
    exists = result.get('output', '').strip() == 'EXISTS'
    info = {'exists': exists, 'is_audio': False, 'size': 0}
    if not exists:
        return info
    file_type_result = env.controller.run_bash_script(f"file -b '{file_path}'", timeout=10)
    file_type = file_type_result.get('output', '').strip()
    is_audio = any((keyword in file_type.lower() for keyword in ['audio', 'mpeg', 'mp3', 'adts']))
    info['is_audio'] = is_audio
    size_result = env.controller.run_bash_script(f"stat -c %s '{file_path}' 2>/dev/null || echo 0", timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
        info['size'] = size
    except:
        info['size'] = 0
    logger.info(f'MP3 file info for {file_path}: {info}')
    return info

def get_desktop_mp4_files__04085b6d(env, config):
    """Get list of MP4 files on Desktop.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        list: List of MP4 filenames on Desktop
    """
    result = env.controller.run_bash_script('cd /home/user/Desktop 2>/dev/null && ls -1 *.mp4 2>/dev/null || echo ""', timeout=10)
    output = result.get('output', '').strip()
    if output:
        return [f.strip() for f in output.split('\n') if f.strip()]
    return []

def get_vlc_snapshot_prefix__bdf11a45(env, config: dict):
    """Get VLC snapshot filename prefix from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        String with snapshot prefix
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    command = f"grep 'snapshot-prefix=' {config_path} 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'snapshot-prefix=' in output:
        prefix = output.split('snapshot-prefix=')[-1].strip()
        return prefix
    else:
        return 'vlc-snap-'

def get_video_file_exists__aec9e92c(env, config):
    """Check if a video file exists at the specified path on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        bool: True if file exists, False otherwise
    """
    path = config.get('path', '')
    result = env.controller.run_bash_script(f'test -f "{path}" && echo "exists" || echo "not_found"', timeout=10)
    output = result.get('output', '').strip()
    return output == 'exists'

def get_mp3_count_in_dir__be089f93(env, config: dict):
    """Count MP3 files in a directory on VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dir_path'

    Returns:
        Number of MP3 files in directory
    """
    dir_path = config['dir_path']
    command = f'find "{dir_path}" -maxdepth 1 -type f -name "*.mp3" | wc -l'
    result = env.controller.run_bash_script(command, timeout=30)
    if result['returncode'] != 0:
        return None
    try:
        count = int(result['output'].strip())
        return count
    except ValueError:
        return None

def get_vlc_always_on_top__b1199d7f2afb5cd10588e7e1e6cdc076(env, config: Dict[str, str]):
    """
    Gets the always-on-top setting from VLC's config file.
    Returns True if always-on-top is enabled, False otherwise.
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\\\AppData\\\\Roaming\\\\vlc\\\\vlcrc'))")['output'].strip()
    else:
        logger.error(f'Unsupported operating system: {os_type}')
        return None
    try:
        content = env.controller.get_file(config_path)
        if not content:
            logger.error('Could not read VLC config file')
            return None
        config_text = content.decode('utf-8')
        always_on_top = None
        for line in config_text.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'video-on-top=' in line:
                always_on_top = line.split('=')[-1].strip()
                break
        if always_on_top is None:
            logger.warning('video-on-top setting not found in VLC config - may not have been configured yet')
            return False
        result = always_on_top == '1'
        logger.info(f'VLC always-on-top setting: {result} (config value: {always_on_top})')
        return result
    except Exception as e:
        logger.error(f'Error getting VLC always-on-top setting: {e}')
        return None

def get_vlc_minimal_view_setting__99af05cec0829d789ac2d7bf7abe0481(env, config: Dict[str, str]):
    """
    Gets the VLC minimal view configuration setting from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' key for output file

    Returns:
        Path to the cached config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_video_info__6e7fd09f(env, config: Dict[str, str]):
    """
    Get video file information including existence, size, and metadata using ffprobe.

    Args:
        env: DesktopEnv instance
        config: Dict with 'path' (output file path), 'source_path' (source file path),
                and 'dest' (local cache destination)

    Returns:
        str: Combined output containing ls, stat, and ffprobe results
    """
    output_path = config.get('path', '/home/user/Desktop/1984_Commercial_Corrected.mp4')
    source_path = config.get('source_path', '/home/user/Desktop/flipped_1984_Apple_Macintosh_Commercial.mp4')
    commands = [f"ls -la '{output_path}'", f"stat -c %s '{output_path}' 2>/dev/null || echo 'ERROR'", f"stat -c %s '{source_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{output_path}' 2>/dev/null || echo 'ERROR'", f"ffprobe -v error -select_streams v:0 -show_entries stream=width,height,rotate -of default=noprint_wrappers=1 '{source_path}' 2>/dev/null || echo 'ERROR'"]
    result_parts = []
    for (i, cmd) in enumerate(commands):
        try:
            output = env.controller.run_bash_script(cmd, timeout=30)
            result_parts.append(output.get('output', ''))
        except Exception as e:
            logger.error(f'Command {i} failed: {e}')
            result_parts.append(f'ERROR: {str(e)}')
    combined_result = result_parts[0] + '---SEPARATOR---' + result_parts[1] + '---FILESIZE_SEP---' + result_parts[2] + '---FFPROBE_OUT---' + result_parts[3] + '---SOURCE_INFO---' + result_parts[4]
    return combined_result

def get_vlc_playback_state__fc6f4242(env, config: Dict[str, str]):
    """
    Gets the current playback state from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        str: Current playback state ('playing', 'paused', 'stopped', etc.), None on error
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            state_element = tree.find('state')
            if state_element is not None and state_element.text:
                state = state_element.text.strip()
                logger.info(f'VLC Playback State: {state}')
                return state
            else:
                logger.warning('State element not found in VLC status XML')
                return None
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC playback state: {e}')
        return None

def get_video_file_properties__4939fb90(env, config):
    """Get video file properties using file command.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: File properties including exists status and file type info
    """
    path = config.get('path', '')
    exists_result = env.controller.run_bash_script(f'test -f "{path}" && echo "1" || echo "0"', timeout=10)
    exists = exists_result.get('output', '').strip() == '1'
    if not exists:
        return {'exists': False, 'file_type': ''}
    file_result = env.controller.run_bash_script(f'file -b "{path}"', timeout=10)
    file_type = file_result.get('output', '').strip()
    return {'exists': True, 'file_type': file_type}

def get_vlc_snapshot_directory__008263c0(env, config: dict):
    """Get VLC snapshot directory from config file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        String with snapshot directory path
    """
    config_path = '/home/user/.config/vlc/vlcrc'
    command = f"grep 'snapshot-path=' {config_path} 2>/dev/null || echo ''"
    result = env.controller.run_bash_script(command, timeout=10)
    output = result.get('output', '').strip()
    if 'snapshot-path=' in output:
        path = output.split('snapshot-path=')[-1].strip()
        return path
    else:
        return '/home/user/Pictures'

def get_audio_extraction_status__d0f84a157daee05fd0a57bb243a6ea5c(env, config: Dict[str, str]):
    """
    Checks if audio file was extracted from video.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'audio_path' key

    Returns:
        dict: Audio file status information
    """
    audio_path = config.get('audio_path', '/home/user/audio.mp3')
    file_check = env.controller.run_bash_script(f"test -f {audio_path} && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = 'EXISTS' in file_check.get('output', '')
    if not file_exists:
        return {'exists': False, 'duration': 0, 'size': 0}
    duration_cmd = f'ffprobe -v error -show_entries format=duration -of default=noprint_wrappers=1:nokey=1 {audio_path}'
    duration_result = env.controller.run_bash_script(duration_cmd, timeout=30)
    try:
        duration = float(duration_result.get('output', '0').strip())
    except:
        duration = 0
    size_cmd = f'stat -c%s {audio_path}'
    size_result = env.controller.run_bash_script(size_cmd, timeout=10)
    try:
        size = int(size_result.get('output', '0').strip())
    except:
        size = 0
    return {'exists': True, 'duration': duration, 'size': size, 'path': audio_path}

def get_audio_file_properties__5d993657(env, config: dict):
    """Get properties of an audio file including size and existence.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Properties including exists, size, path
    """
    file_path = config.get('path', '')
    result = env.controller.run_bash_script(f"if [ -f '{file_path}' ]; then stat -c '%s' '{file_path}'; else echo '0'; fi", timeout=10)
    file_size = 0
    exists = False
    if result['returncode'] == 0:
        output = result['output'].strip()
        if output and output != '0':
            try:
                file_size = int(output)
                exists = file_size > 0
            except ValueError:
                pass
    return {'exists': exists, 'size': file_size, 'path': file_path}

def get_vlc_loop_status__9e0433d6(env, config: Dict[str, str]):
    """
    Gets the loop/repeat status from VLC's HTTP interface.

    Args:
        env: DesktopEnv instance
        config: Configuration dict

    Returns:
        bool: True if loop is enabled, False otherwise, None on error
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/status.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=10)
        if response.status_code == 200:
            tree = ElementTree.fromstring(response.content)
            loop_element = tree.find('loop')
            if loop_element is not None and loop_element.text:
                loop_status = loop_element.text.strip().lower() in ['true', '1', 'yes']
                logger.info(f'VLC Loop Status: {loop_status}')
                return loop_status
            else:
                logger.info('Loop element not found in status XML, assuming False')
                return False
        else:
            logger.error(f'Failed to get VLC status. Status code: {response.status_code}')
            return None
    except Exception as e:
        logger.error(f'Error getting VLC loop status: {e}')
        return None

def get_video_metadata_file__bc3a25a7baab15992708f46f1d51e584(env, config: Dict[str, str]):
    """
    Checks if video metadata was extracted to a file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'metadata_path' key

    Returns:
        dict: Metadata file status information
    """
    metadata_path = config.get('metadata_path', '/home/user/metadata.txt')
    file_check = env.controller.run_bash_script(f"test -f {metadata_path} && echo 'EXISTS' || echo 'NOT_EXISTS'", timeout=10)
    file_exists = 'EXISTS' in file_check.get('output', '')
    if not file_exists:
        return {'exists': False, 'content': '', 'has_duration': False, 'has_resolution': False, 'duration_value': None, 'resolution_value': None}
    content_result = env.controller.run_bash_script(f'cat {metadata_path}', timeout=10)
    content = content_result.get('output', '')
    content_lower = content.lower()
    error_keywords = ['unknown', 'failed', 'error', 'not found', 'n/a']
    has_errors = any((keyword in content_lower for keyword in error_keywords))
    duration_pattern = '(\\d+:\\d+:\\d+|\\d+:\\d+|\\d+\\s*(min|sec|second|hour|hr)s?)'
    duration_match = re.search(duration_pattern, content_lower)
    has_duration_keyword = 'duration' in content_lower or 'length' in content_lower or 'time' in content_lower
    has_duration = duration_match is not None and has_duration_keyword and (not has_errors)
    duration_value = duration_match.group(0) if duration_match else None
    resolution_pattern = '\\d+\\s*x\\s*\\d+'
    resolution_match = re.search(resolution_pattern, content_lower)
    width_height_pattern = 'width[:\\s]+\\d+.*height[:\\s]+\\d+|height[:\\s]+\\d+.*width[:\\s]+\\d+'
    width_height_match = re.search(width_height_pattern, content_lower)
    has_resolution_keyword = 'resolution' in content_lower or ('width' in content_lower and 'height' in content_lower)
    has_resolution = (resolution_match is not None or width_height_match is not None) and has_resolution_keyword and (not has_errors)
    resolution_value = resolution_match.group(0) if resolution_match else width_height_match.group(0) if width_height_match else None
    return {'exists': True, 'content': content, 'has_duration': has_duration, 'has_resolution': has_resolution, 'duration_value': duration_value, 'resolution_value': resolution_value, 'has_errors': has_errors, 'path': metadata_path}

def get_vlc_recordings_folder__055ce0bbe52cb194135b41b067eca986(env, config: Dict[str, str]):
    """
    Gets the VLC recording folder path from configuration file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'dest' key for output file

    Returns:
        Path to the cached config file
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/.config/vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Darwin':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~/Library/Preferences/org.videolan.vlc/vlcrc'))")['output'].strip()
    elif os_type == 'Windows':
        config_path = env.controller.execute_python_command("import os; print(os.path.expanduser('~\\AppData\\Roaming\\vlc\\vlcrc'))")['output'].strip()
    else:
        raise Exception('Unsupported operating system', os_type)
    _path = os.path.join(env.cache_dir, config['dest'])
    content = env.controller.get_file(config_path)
    with open(_path, 'wb') as f:
        f.write(content)
    return _path

def get_vlc_playlist_count__fc73aaff(env, config: dict):
    """
    Get the number of items in VLC's current playlist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters (not used but required)

    Returns:
        int: Number of items in playlist
    """
    host = env.vm_ip
    port = env.vlc_port
    password = 'password'
    url = f'http://{host}:{port}/requests/playlist.xml'
    try:
        response = requests.get(url, auth=('', password), timeout=5)
        if response.status_code == 200:
            root = ET.fromstring(response.content)
            items = root.findall('.//leaf')
            return len(items)
    except Exception as e:
        logger.error(f'Failed to get VLC playlist: {e}')
    return 0

def get_srt_and_video_status__ed96ceb6(env, config: dict):
    """
    Get subtitle file and check if video still has embedded subtitles.

    Args:
        env: Desktop environment
        config: Configuration dict

    Returns:
        dict: Contains srt_path and has_subtitles status
    """
    srt_vm_path = config['srt_path']
    srt_dest = config['srt_dest']
    video_path = config.get('video_path', '/home/user/video.mp4')
    result = {'srt_path': None, 'has_subtitles': True}
    cache_path = os.path.join(env.cache_dir, srt_dest)
    file_content = env.controller.get_file(srt_vm_path)
    if file_content is not None:
        with open(cache_path, 'wb') as f:
            f.write(file_content)
        result['srt_path'] = cache_path
    cmd = f'ffmpeg -i {video_path}'
    ffmpeg_result = env.controller.run_bash_script(cmd, timeout=30)
    output = ffmpeg_result.get('error', '') + ffmpeg_result.get('output', '')
    if 'Stream #0:2' not in output or 'Subtitle' not in output:
        result['has_subtitles'] = False
    return result

def get_mp3_audio_info__5e12822f6a68285e2467088dd7d25598(env, config):
    """
    Get comprehensive MP3 audio file information including format verification,
    audio properties, and source video comparison.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for MP3 file and optional 'source_video_path'

    Returns:
        dict: Audio information including:
            - exists: bool - whether file exists
            - file_type: str - file type from 'file' command
            - is_mp3: bool - whether file is MP3 format
            - size_mb: float - file size in MB
            - has_audio: bool - whether file contains valid audio stream
            - codec: str - audio codec name
            - duration: float - audio duration in seconds
            - sample_rate: int - audio sample rate in Hz
            - bit_rate: int - audio bit rate in bps
            - source_duration: float - source video duration (if source_video_path provided)
            - duration_match: bool - whether durations match (within tolerance)
    """
    mp3_path = config.get('path', '')
    source_video_path = config.get('source_video_path', '')
    result = {'exists': False, 'file_type': '', 'is_mp3': False, 'size_mb': 0.0, 'has_audio': False, 'codec': '', 'duration': 0.0, 'sample_rate': 0, 'bit_rate': 0, 'source_duration': 0.0, 'duration_match': False}
    check_exists = env.controller.run_bash_script(f"if [ -f '{mp3_path}' ]; then echo 'exists'; else echo 'not_found'; fi", timeout=10)
    if 'not_found' in check_exists.get('output', ''):
        logger.warning(f'MP3 file not found: {mp3_path}')
        return result
    result['exists'] = True
    file_type_result = env.controller.run_bash_script(f"file -b '{mp3_path}'", timeout=10)
    result['file_type'] = file_type_result.get('output', '').strip()
    file_type_lower = result['file_type'].lower()
    result['is_mp3'] = 'mpeg' in file_type_lower and 'audio' in file_type_lower
    size_result = env.controller.run_bash_script(f"stat -c %s '{mp3_path}'", timeout=10)
    try:
        size_bytes = int(size_result.get('output', '0').strip())
        result['size_mb'] = round(size_bytes / (1024 * 1024), 2)
    except:
        result['size_mb'] = 0.0
    ffprobe_cmd = f"ffprobe -v error -show_entries stream=codec_name,codec_type,sample_rate,bit_rate -show_entries format=duration -of default=noprint_wrappers=1 '{mp3_path}'"
    ffprobe_result = env.controller.run_bash_script(ffprobe_cmd, timeout=30)
    ffprobe_output = ffprobe_result.get('output', '')
    if ffprobe_output:
        lines = ffprobe_output.strip().split('\n')
        for line in lines:
            if '=' in line:
                (key, value) = line.split('=', 1)
                key = key.strip()
                value = value.strip()
                if key == 'codec_type' and value == 'audio':
                    result['has_audio'] = True
                elif key == 'codec_name':
                    result['codec'] = value
                elif key == 'duration':
                    try:
                        result['duration'] = float(value)
                    except:
                        pass
                elif key == 'sample_rate':
                    try:
                        result['sample_rate'] = int(value)
                    except:
                        pass
                elif key == 'bit_rate':
                    try:
                        result['bit_rate'] = int(value)
                    except:
                        pass
    if source_video_path:
        source_ffprobe_cmd = f"ffprobe -v error -show_entries format=duration -of default=noprint_wrappers=1:nokey=1 '{source_video_path}'"
        source_result = env.controller.run_bash_script(source_ffprobe_cmd, timeout=30)
        source_output = source_result.get('output', '').strip()
        try:
            result['source_duration'] = float(source_output)
            if result['duration'] > 0 and result['source_duration'] > 0:
                tolerance = max(2.0, result['source_duration'] * 0.05)
                duration_diff = abs(result['duration'] - result['source_duration'])
                result['duration_match'] = duration_diff <= tolerance
        except:
            pass
    logger.info(f'MP3 audio info for {mp3_path}: {result}')
    return result
