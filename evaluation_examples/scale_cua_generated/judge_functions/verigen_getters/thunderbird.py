"""VeriGen generated judge functions.

Source: getters.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, time
from datetime import datetime, timedelta
from datetime import time
from desktop_env.evaluators.getters.chrome import GoogleAuth, GoogleDrive
from desktop_env.evaluators.getters.chrome import get_bookmarks
from desktop_env.evaluators.getters.chrome import get_open_tabs_info, get_bookmarks
from desktop_env.evaluators.getters.file import get_vm_file
from desktop_env.evaluators.getters.general import get_vm_command_line
from desktop_env.evaluators.metrics.utils import read_cell_value
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_LINE_SPACING
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.oxml.ns import qn
from docx.oxml.shape import CT_Picture
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.shared import Pt
from docx.shared import RGBColor
from docx.table import Table
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from email import message_from_string
from email import policy
from email.utils import parseaddr
from email.utils import parseaddr, parsedate_to_datetime
from email.utils import parsedate_to_datetime
from io import BytesIO
from io import StringIO
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from lxml.etree import _Element as Element
from odf import teletype
from odf import text, style, teletype
from odf import text, teletype
from odf.opendocument import load
from odf.style import TextProperties
from odf.text import Span, P
from openpyxl.comments import Comment
from openpyxl.styles import Alignment
from openpyxl.styles import Color
from openpyxl.styles import Font
from openpyxl.styles import PatternFill
from openpyxl.styles.fills import PatternFill
from openpyxl.utils import get_column_letter
from openpyxl.utils import get_column_letter, column_index_from_string
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
from pydrive.files import GoogleDriveFile, GoogleDriveFileList
from pypdf import PdfReader
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, Optional
from typing import Any, Dict, Tuple
from typing import Any, Optional
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, Iterable
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Optional, Tuple
from typing import Dict, Any, List, Set
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Optional, Union
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, List
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, Any, Union, List
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Any, Optional
from typing import Dict, List, Set
from typing import Dict, List, Tuple
from typing import Dict, Optional
from typing import Dict, Set
from typing import Dict, Tuple, Optional
from typing import List
from typing import List, Any
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Optional
from typing import List, Tuple
from typing import Optional
from typing import Optional, Dict, Any
from typing import Optional, Dict, List, Any
from typing import Optional, Tuple
from typing import TypeVar, Dict
from urllib.parse import unquote
from urllib.parse import unquote, urlparse, parse_qs
from urllib.parse import urlparse, parse_qs, unquote
from urllib.parse import urlparse, unquote
from xml import etree
from xml.etree import ElementTree
from xml.etree import ElementTree as ET
from zipfile import ZipFile
import PyPDF2
import ast
import calendar
import csv
import cv2
import email
import fitz
import glob
import gzip
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import logging, tempfile, os
import lxml.etree
import mailbox
import numpy as np
import openpyxl
import openpyxl, tempfile, os as _os
import os
import os.path
import platform
import pytz
import random
import re
import requests
import shutil
import sqlite3
import struct
import subprocess
import sys
import tarfile
import tempfile
import time
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['get_thunderbird_prefs__0d04b8a7', 'get_first_contact_email__105b8e17', 'get_thunderbird_prefs__f50e981c', 'get_thunderbird_prefs__625e2b37', 'get_thunderbird_recipient_count__c9ce3f52', 'get_thunderbird_prefs__dced30f9', 'get_thunderbird_attachment_state__d38192b0', 'get_thunderbird_prefs__f0f94c7477cc672daff3cf5ec91ccac4', 'get_thunderbird_prefs__7034ec46', 'get_thunderbird_prefs__1f5ab547a29208fe81de884f80a1716d', 'get_thunderbird_all_addresses__d8fe40b3', 'get_thunderbird_digest_subjects__5a1ab509', 'get_tb_first3_with_dates__843811624c18ed6ed65e7d4d877f3122', 'get_csv_contact_emails__0fa83097089d9008d43d2c029fbbd194', 'get_thunderbird_prefs__249498a7cd34d43ebd7ef238b7ba6bee', 'get_thunderbird_draft_attachment__d38192b0_aug9', 'get_tb_sender_counts__6a3e0f1dfc20fc63f235dcd50f4dfaaf', 'get_thunderbird_folder_and_star_check__7e32e736', 'get_thunderbird_unique_senders__0399616e', 'get_thunderbird_archive_status__5fbabc93', 'get_thunderbird_to_recipients__c9ce3f52', 'get_thunderbird_folder_counts__856f6f3499821b31f240098182dcc479', 'get_email_count_file__fcf6b1e44de70adba5822f1fa1e262b7', 'get_thunderbird_two_folders_exist__386310fb', 'get_tb_email_summary__095c028a04cb1496af16d82dbb5f9c21', 'get_email_formatting__fba523fac66398bfd72f79ad1e012e07', 'get_thunderbird_and_gdrive_pdfs__69a89619e7273cf1d986af62cad42166', 'get_thunderbird_subject_date__2284321c', 'get_thunderbird_cc_emails__f218f3c2', 'get_thunderbird_folder_emails__4bba48e3', 'get_thunderbird_domain_count__1744c07b', 'get_thunderbird_oldest_three__0c7bac1c', 'get_thunderbird_txt_backup__dd834e2c641c360491704c9db4c6f362', 'get_thunderbird_prefs__616461c5ba008feaf990bbd287063304', 'get_thunderbird_prefs__a3ecd62701cd6d8575ca51c05b0fdb89', 'get_thunderbird_date_from_subject__a92c1922', 'get_thunderbird_prefs__51cc10d9d3791de4bc35e6d331e3ef94', 'get_thunderbird_prefs__5b041e31e4566087bb3bbe0b347cd7ff', 'get_thunderbird_prefs__bd92e86076085c3fd0cc432ab151ca8f', 'get_thunderbird_prefs__cd29a448', 'get_tb_filtered_emails__01abb20597fbc4723613687c7c60c0d3', 'get_default_email_client__b929c586', 'get_email_count_file__c95564e31fba18e60360c502646bdd5f', 'get_thunderbird_folder_and_email_count__a91026e2', 'get_thunderbird_prefs__17405a90', 'get_thunderbird_sender_count__44f17dcf', 'get_thunderbird_prefs__4d8c063289348a4e74b61401db33c326', 'get_thunderbird_folder_counts__4f3b75dbef7eff0e54cd50b24ff5fb5e', 'get_txt_contains_email__a3b9c6d3', 'get_thunderbird_and_chrome_state__0f025d55a86976e1ae415a210c74f719', 'get_email_attachments__d38192b0', 'get_test_emails_backup__a9ff16fc1be9331aa02283bf7c168b3e', 'get_thunderbird_prefs__550830d3', 'get_thunderbird_prefs__e06bf170da6e0042eda0ceb7f1f9f833', 'get_thunderbird_prefs__6e76242e9ee0066f62efb661397a7719', 'get_thunderbird_prefs__dddfffe093b014a09b5e9ec4efaf461f', 'get_thunderbird_email_read_status__a63fb94e', 'get_email_subjects_file__df60c313c1f3b2e712e9756c47386ddc', 'get_thunderbird_subject__c9ce3f52', 'get_default_email_client__119177bf', 'get_thunderbird_email_count_and_folder__3c15b6d0', 'get_thunderbird_filters__f57034f8', 'get_thunderbird_filter_rules__8a2db6ad', 'get_thunderbird_three_cols__07170c7f', 'get_thunderbird_prefs__734b4908962171d93d8cfb6878a651d2', 'get_thunderbird_prefs__497f2c49']

def get_thunderbird_prefs__0d04b8a7(env, config):
    """Get Thunderbird prefs.js file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_first_contact_email__105b8e17(env, config):
    """
    Get the email from the first contact row in XLSX.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with xlsx_path

    Returns:
        str: Email address from first contact, or empty string if not found
    """
    xlsx_path = config.get('xlsx_path', '/home/user/Desktop/contacts.xlsx')
    xlsx_bytes = env.controller.get_file(xlsx_path)
    if not xlsx_bytes:
        return ''
    try:
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(xlsx_bytes)
            tmp_path = tmp.name
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        email = ws.cell(row=2, column=5).value or ''
        os.unlink(tmp_path)
        return email
    except Exception:
        return ''

def get_thunderbird_prefs__f50e981c(env, config):
    """Get Thunderbird prefs.js file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_prefs__625e2b37(env, config):
    """
    Get Thunderbird prefs.js file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'profile_path' key

    Returns:
        str: Path to the downloaded prefs.js file
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_recipient_count__c9ce3f52(env, config):
    """Get the list of recipient names in the To field.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (unused but required by framework)

    Returns:
        list: List of recipient names/identifiers
    """
    accessibility_tree = env.controller.get_accessibility_tree()
    if not accessibility_tree:
        return []
    recipients = []
    lines = accessibility_tree.split('\n')
    for line in lines:
        if 'MsgHeadersToolbar' in line and 'address-pill' in line and ('pill-label' in line):
            if 'pill-label:' in line:
                parts = line.split('pill-label:')
                if len(parts) > 1:
                    label_part = parts[1].strip()
                    if label_part.startswith('"'):
                        end_quote = label_part.find('"', 1)
                        if end_quote > 0:
                            name = label_part[1:end_quote].strip()
                            recipients.append(name)
    return recipients

def get_thunderbird_prefs__dced30f9(env, config):
    """
    Get Thunderbird prefs.js file from the VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'profile_path' key

    Returns:
        str: Path to the downloaded prefs.js file
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_attachment_state__d38192b0(env, config):
    """
    Get the state of Thunderbird compose window including attachment and window status.

    This function verifies:
    1. That the Thunderbird compose window is still open (not closed or sent)
    2. That the specified attachment is present in the compose window

    Args:
        env: Environment object
        config: Configuration dict with 'subject' and 'attachment_name' keys

    Returns:
        dict: {
            'window_open': bool - Whether the compose window is still open
            'attachment_present': bool - Whether the attachment is in the window
            'subject': str - The email subject being checked
            'attachment_name': str - The attachment filename being checked
        }
    """
    subject = config.get('subject', 'New-month AWS Bill')
    attachment_name = config.get('attachment_name', 'aws-bill.pdf')
    writer_window = get_thunderbird_writer_at(subject)
    window_open = writer_window is not None
    attachment_present = False
    if window_open:
        attachment_present = check_attachment(subject, attachment_name)
    return {'window_open': window_open, 'attachment_present': attachment_present, 'subject': subject, 'attachment_name': attachment_name}

def get_thunderbird_prefs__f0f94c7477cc672daff3cf5ec91ccac4(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get Thunderbird preferences file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        Path to the downloaded prefs.js file in cache directory, or None if download fails
    """
    path = config['path']
    dest = config.get('dest', 'thunder-prefs.js')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_prefs__7034ec46(env, config: dict):
    """Get Thunderbird preferences file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'prefs_path'

    Returns:
        str: Path to downloaded prefs.js file
    """
    prefs_path = config.get('prefs_path', '/home/user/.thunderbird/t5q2a5hp.default-release/prefs.js')
    result = env.controller.get_file(prefs_path)
    if result is None:
        return None
    cache_path = os.path.join(env.cache_dir, 'prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(result)
    return cache_path

def get_thunderbird_prefs__1f5ab547a29208fe81de884f80a1716d(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get Thunderbird preferences file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        Path to the downloaded prefs.js file in cache directory, or None if download fails
    """
    path = config['path']
    dest = config.get('dest', 'thunder-prefs.js')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_all_addresses__d8fe40b3(env, config: dict):
    """Extract sender and CC email addresses from Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of tuples (sender_address, cc_address)
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        data = []
        for row_idx in range(2, ws.max_row + 1):
            sender_addr = ws.cell(row_idx, 1).value
            cc_addr = ws.cell(row_idx, 2).value
            data.append((sender_addr, cc_addr))
        return data
    finally:
        os.unlink(tmp_path)

def get_thunderbird_digest_subjects__5a1ab509(env, config: dict):
    """Extract subjects containing 'Digest' from Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of subject lines containing 'Digest'
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        for row_idx in range(1, ws.max_row + 1):
            for col_idx in range(2, ws.max_column + 1):
                if ws.cell(row_idx, col_idx).value is not None:
                    return []
        subjects = []
        start_row = 1
        first_cell = ws.cell(1, 1).value
        if first_cell and 'Digest' not in str(first_cell):
            start_row = 2
        for row_idx in range(start_row, ws.max_row + 1):
            subject = ws.cell(row_idx, 1).value
            if subject:
                subject_str = str(subject).strip()
                if 'digest' in subject_str.lower():
                    subjects.append(subject_str)
        return subjects
    finally:
        os.unlink(tmp_path)

def get_tb_first3_with_dates__843811624c18ed6ed65e7d4d877f3122(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Extract first 3 emails with date information from spreadsheet.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of email dictionaries with date info
    """
    file_path = config.get('path', '')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        emails = []
        row_count = 0
        for (row_idx, row) in enumerate(ws.iter_rows(min_row=1, values_only=True), start=1):
            if row_idx == 1 and row and isinstance(row[0], str):
                if any((keyword in str(row[0]).lower() for keyword in ['sender', 'subject', 'date'])):
                    continue
            if row and any((cell is not None and str(cell).strip() for cell in row)):
                email_data = {'sender_name': str(row[0]).strip() if row[0] is not None else '', 'subject': str(row[1]).strip() if len(row) > 1 and row[1] is not None else '', 'date': str(row[2]).strip() if len(row) > 2 and row[2] is not None else ''}
                emails.append(email_data)
                row_count += 1
                if row_count >= 3:
                    break
        wb.close()
        return emails
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_csv_contact_emails__0fa83097089d9008d43d2c029fbbd194(env, config):
    """Get list of primary email addresses from CSV file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of email addresses
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.csv', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        emails = []
        with open(tmp_path, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            for row in reader:
                email = row.get('Primary Email', '').strip()
                if email:
                    emails.append(email)
        return emails
    finally:
        os.unlink(tmp_path)

def get_thunderbird_prefs__249498a7cd34d43ebd7ef238b7ba6bee(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get Thunderbird preferences file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        Path to the downloaded prefs.js file in cache directory, or None if download fails
    """
    path = config['path']
    dest = config.get('dest', 'thunder-prefs.js')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_draft_attachment__d38192b0_aug9(env, config):
    """
    Check Thunderbird email draft for correct attachment and verify email is not sent.

    This getter verifies:
    1. Email exists in Drafts folder (not sent)
    2. Email has the expected attachment (.aws-bill-mail-body.html)
    3. Attachment file size is reasonable (not empty)

    Args:
        env: Environment object
        config: Configuration dict with:
            - subject: Email subject to search for
            - attachment_name: Expected attachment filename
            - thunderbird_profile: Path to Thunderbird profile (optional)

    Returns:
        dict: {
            'draft_exists': bool,
            'not_sent': bool,
            'has_attachment': bool,
            'attachment_name': str or None,
            'attachment_size': int or None
        }
    """
    subject = config.get('subject', 'New-month AWS Bill')
    attachment_name = config.get('attachment_name', '.aws-bill-mail-body.html')
    profile_path = config.get('thunderbird_profile', '/home/user/.thunderbird/t5q2a5hp.default-release')
    result = {'draft_exists': False, 'not_sent': False, 'has_attachment': False, 'attachment_name': None, 'attachment_size': None}
    drafts_path = os.path.join(profile_path, 'Mail/Local Folders/Drafts')
    sent_path = os.path.join(profile_path, 'Mail/Local Folders/Sent')
    try:
        drafts_content = env.controller.get_file(drafts_path)
        if drafts_content:
            drafts_text = drafts_content.decode('utf-8', errors='ignore')
            if f'Subject: {subject}' in drafts_text:
                result['draft_exists'] = True
                logger.info(f"Found email with subject '{subject}' in Drafts")
                messages = drafts_text.split('From - ')
                for msg in messages:
                    if f'Subject: {subject}' in msg:
                        attachment_pattern = 'Content-Disposition:\\s*attachment;\\s*filename="([^"]+)"'
                        matches = re.findall(attachment_pattern, msg, re.IGNORECASE)
                        if matches:
                            for match in matches:
                                if attachment_name in match:
                                    result['has_attachment'] = True
                                    result['attachment_name'] = match
                                    logger.info(f'Found attachment: {match}')
                                    attachment_section_pattern = f'filename="{re.escape(match)}".*?(?=Content-Type:|From - |\\Z)'
                                    attachment_section = re.search(attachment_section_pattern, msg, re.DOTALL | re.IGNORECASE)
                                    if attachment_section:
                                        attachment_data = attachment_section.group(0)
                                        result['attachment_size'] = len(attachment_data)
                                        if len(attachment_data) > 500:
                                            logger.info(f'Attachment appears to have content (size estimate: {len(attachment_data)} bytes)')
                                    break
                        break
            else:
                logger.info(f"Email with subject '{subject}' not found in Drafts")
        else:
            logger.warning(f'Could not read Drafts file at {drafts_path}')
    except Exception as e:
        logger.warning(f'Error checking Drafts folder: {e}')
    try:
        sent_content = env.controller.get_file(sent_path)
        if sent_content:
            sent_text = sent_content.decode('utf-8', errors='ignore')
            if f'Subject: {subject}' not in sent_text:
                result['not_sent'] = True
                logger.info(f'Confirmed email is not in Sent folder')
            else:
                logger.warning(f"Email with subject '{subject}' found in Sent folder - task requirement violated!")
        else:
            result['not_sent'] = True
            logger.info(f'Sent folder is empty - email not sent')
    except Exception as e:
        logger.warning(f'Error checking Sent folder: {e}')
        if result['draft_exists']:
            result['not_sent'] = True
    try:
        attachment_file_path = config.get('attachment_file_path', f'/home/user/{attachment_name}')
        file_content = env.controller.get_file(attachment_file_path)
        if file_content and len(file_content) > 0:
            logger.info(f'Verified source attachment file exists at {attachment_file_path} with size {len(file_content)} bytes')
    except Exception as e:
        logger.debug(f'Could not verify source attachment file: {e}')
    logger.info(f'Draft check result: {result}')
    return result

def get_tb_sender_counts__6a3e0f1dfc20fc63f235dcd50f4dfaaf(env, config: Dict[str, Any]) -> Dict[str, int]:
    """
    Extract sender email counts from a LibreOffice Calc spreadsheet.
    Expected format: Column A contains sender names, Column B contains counts.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict mapping sender names to counts
    """
    file_path = config.get('path', '')
    if not file_path:
        return {}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        sender_counts = {}
        for (row_idx, row) in enumerate(ws.iter_rows(min_row=1, values_only=True), start=1):
            if row_idx == 1 and row and isinstance(row[0], str):
                if any((keyword in str(row[0]).lower() for keyword in ['sender', 'name', 'from'])):
                    continue
            if row and row[0] is not None:
                sender = str(row[0]).strip()
                count = int(row[1]) if len(row) > 1 and row[1] is not None else 0
                if sender:
                    sender_counts[sender] = count
        wb.close()
        return sender_counts
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_thunderbird_folder_and_star_check__7e32e736(env, config: dict):
    """Check if a Thunderbird folder exists and if Bills folder has starred emails.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - profile_path: Path to Thunderbird profile
            - folder_to_check: Name of folder to check existence
            - bills_folder_path: Path to Bills folder file

    Returns:
        dict: {
            "folder_exists": bool,
            "has_starred_email": bool
        }
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    folder_to_check = config.get('folder_to_check', 'Important')
    bills_folder_path = config.get('bills_folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills')
    result = {'folder_exists': False, 'has_starred_email': False}
    folder_path = os.path.join(profile_path, 'Mail', 'Local Folders', folder_to_check)
    folder_file_bytes = env.controller.get_file(folder_path)
    if folder_file_bytes is not None:
        result['folder_exists'] = True
        logger.info(f'Folder {folder_to_check} exists')
    else:
        logger.info(f'Folder {folder_to_check} does not exist')
    bills_bytes = env.controller.get_file(bills_folder_path)
    if bills_bytes is not None:
        bills_content = bills_bytes.decode('utf-8', errors='ignore')
        if re.search('X-Mozilla-Status: 000[23]', bills_content):
            result['has_starred_email'] = True
            logger.info('Found starred email in Bills folder')
        else:
            logger.info('No starred email found in Bills folder')
    else:
        logger.info('Bills folder not found')
    return result

def get_thunderbird_unique_senders__0399616e(env, config: dict):
    """Extract unique sender email addresses from Excel report and verify Thunderbird state.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Contains 'senders' (sorted list), 'thunderbird_verified' (bool),
              'email_count_verified' (bool), or None if file doesn't exist
    """
    import tempfile
    import os
    import glob
    result = {'senders': None, 'thunderbird_verified': False, 'email_count_verified': False}
    try:
        thunderbird_profile_base = '/home/user/.thunderbird'
        daily_folder_patterns = [f'{thunderbird_profile_base}/*/ImapMail/*/daily*', f'{thunderbird_profile_base}/*/Mail/*/daily*']
        found_daily_folder = False
        email_file_path = None
        for pattern in daily_folder_patterns:
            matches = glob.glob(pattern)
            if matches:
                for match in matches:
                    if not match.endswith('.msf'):
                        found_daily_folder = True
                        email_file_path = match
                        break
            if found_daily_folder:
                break
        if found_daily_folder and email_file_path:
            result['thunderbird_verified'] = True
            try:
                with open(email_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    email_count = content.count('\nFrom - ')
                    if content.startswith('From - '):
                        email_count += 1
                    if email_count >= 5:
                        result['email_count_verified'] = True
            except:
                pass
    except:
        pass
    try:
        file_bytes = env.controller.get_file(config['path'])
    except (FileNotFoundError, Exception) as e:
        return None
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        senders = set()
        for row_idx in range(1, ws.max_row + 1):
            sender_email = ws.cell(row_idx, 1).value
            if sender_email and isinstance(sender_email, str):
                if '@' in sender_email:
                    senders.add(sender_email)
        result['senders'] = sorted(list(senders))
        return result
    except Exception as e:
        return None
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_thunderbird_archive_status__5fbabc93(env, config: dict):
    """Check if the 'Paper Recommendation' email was archived.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'archive_path' and 'email_subject'

    Returns:
        dict: Contains 'found' (bool) and 'details' (str) about the search result
    """
    archive_path = config.get('archive_path', '/home/user/.thunderbird/t5q2a5hp.default-release/ImapMail/outlook.office365.com/Archives')
    email_subject = config.get('email_subject', 'Paper Recommendation')
    try:
        archive_file = env.controller.get_file(archive_path)
        if archive_file is None:
            logger.warning(f'Archive file not found at {archive_path}')
            return {'found': False, 'details': 'Archive file does not exist'}
        if len(archive_file) == 0:
            logger.info('Archive file is empty')
            return {'found': False, 'details': 'Archive file is empty'}
        try:
            archive_content = archive_file.decode('utf-8', errors='ignore')
        except Exception as e:
            logger.error(f'Error decoding archive file: {e}')
            return {'found': False, 'details': f'Error decoding file: {str(e)}'}
        subject_patterns = [f'Subject: {email_subject}', f'Subject: Re: {email_subject}', f'Subject: Fwd: {email_subject}', email_subject]
        found = any((pattern in archive_content for pattern in subject_patterns))
        if found:
            logger.info(f"Found email with subject '{email_subject}' in archive")
            return {'found': True, 'details': f"Email '{email_subject}' found in archive"}
        else:
            logger.warning(f"Email with subject '{email_subject}' not found in archive")
            return {'found': False, 'details': f'Email not found. Archive size: {len(archive_file)} bytes'}
    except Exception as e:
        logger.error(f'Error checking archive status: {e}')
        return {'found': False, 'details': f'Error: {str(e)}'}

def get_thunderbird_to_recipients__c9ce3f52(env, config):
    """Get all email addresses in the To field.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (unused but required by framework)

    Returns:
        list: List of email addresses in To field
    """
    accessibility_tree = env.controller.get_accessibility_tree()
    if not accessibility_tree:
        return []
    recipients = []
    lines = accessibility_tree.split('\n')
    for line in lines:
        if 'MsgHeadersToolbar' in line and 'To' in line and ('pill-label' in line) and ('name=' in line):
            parts = line.split('name=')
            if len(parts) > 1:
                name_part = parts[-1]
                if '"' in name_part:
                    email_text = name_part.split('"')[1]
                    if '<' in email_text and '>' in email_text:
                        email = email_text.split('<')[1].split('>')[0]
                    else:
                        email = email_text
                    recipients.append(email.strip())
    return recipients

def get_thunderbird_folder_counts__856f6f3499821b31f240098182dcc479(env, config):
    """Get message counts and metadata from multiple Thunderbird folders.

    Args:
        env: DesktopEnv instance
        config: Dict with 'paths' list of folder file paths

    Returns:
        dict: Contains 'counts' (folder -> count), 'bills_most_recent' (subject of most recent message in Bills),
              and 'have_seen_messages' (list of message subjects in have_seen folder)
    """
    paths = config.get('paths', [])
    if not paths:
        logger.error('No paths specified in config')
        return None
    counts = {}
    bills_messages = []
    have_seen_messages = []
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if not file_bytes:
                logger.info(f'File not found: {path}')
                counts[path] = 0
                continue
            content = file_bytes.decode('utf-8', errors='ignore')
            messages = content.split('From - ')
            count = 0
            folder_messages = []
            for msg in messages[1:]:
                if not msg.strip():
                    continue
                if re.search('X-Mozilla-Status: 000[89]', msg):
                    continue
                count += 1
                try:
                    parsed_msg = message_from_string(msg)
                    subject = parsed_msg.get('Subject', '').strip()
                    date_str = parsed_msg.get('Date', '')
                    try:
                        date = parsedate_to_datetime(date_str)
                    except:
                        date = None
                    folder_messages.append({'subject': subject, 'date': date, 'date_str': date_str})
                except Exception as e:
                    logger.warning(f'Failed to parse message in {path}: {e}')
            counts[path] = count
            logger.info(f'Found {count} messages in {path}')
            if 'Bills' in path and 'have_seen' not in path:
                bills_messages = folder_messages
            elif 'have_seen' in path:
                have_seen_messages = folder_messages
        except Exception as e:
            logger.error(f'Error reading {path}: {e}')
            counts[path] = None
    bills_most_recent = None
    if bills_messages:
        messages_with_dates = [m for m in bills_messages if m['date'] is not None]
        if messages_with_dates:
            messages_with_dates.sort(key=lambda m: m['date'], reverse=True)
            bills_most_recent = messages_with_dates[0]['subject']
            logger.info(f'Most recent message in Bills: {bills_most_recent}')
        else:
            bills_most_recent = bills_messages[0]['subject'] if bills_messages else None
            logger.warning('No valid dates found in Bills messages, using first message')
    have_seen_subjects = [m['subject'] for m in have_seen_messages]
    result = {'counts': counts, 'bills_most_recent': bills_most_recent, 'have_seen_subjects': have_seen_subjects}
    logger.info(f'Result: counts={counts}, bills_most_recent={bills_most_recent}, have_seen_subjects={have_seen_subjects}')
    return result

def get_email_count_file__fcf6b1e44de70adba5822f1fa1e262b7(env, config):
    """
    Get the content of email count file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to count file

    Returns:
        str: Content of the count file (number as string)
    """
    file_path = config.get('path', '/home/user/inbox_count.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    content = file_bytes.decode('utf-8').strip()
    return content

def get_thunderbird_two_folders_exist__386310fb(env, config: dict):
    """Check if two Thunderbird folders exist.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - profile_path: Path to Thunderbird profile
            - folders: List of folder names to check

    Returns:
        dict: {
            "personal_exists": bool,
            "work_exists": bool
        }
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    folders = config.get('folders', ['Personal', 'Work'])
    result = {'personal_exists': False, 'work_exists': False}
    personal_path = os.path.join(profile_path, 'Mail', 'Local Folders', folders[0])
    personal_bytes = env.controller.get_file(personal_path)
    if personal_bytes is not None:
        result['personal_exists'] = True
        logger.info(f'Folder {folders[0]} exists')
    else:
        logger.info(f'Folder {folders[0]} does not exist')
    work_path = os.path.join(profile_path, 'Mail', 'Local Folders', folders[1])
    work_bytes = env.controller.get_file(work_path)
    if work_bytes is not None:
        result['work_exists'] = True
        logger.info(f'Folder {folders[1]} exists')
    else:
        logger.info(f'Folder {folders[1]} does not exist')
    return result

def get_tb_email_summary__095c028a04cb1496af16d82dbb5f9c21(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract summary statistics from email report spreadsheet.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        Dict containing summary statistics
    """
    file_path = config.get('path', '')
    if not file_path:
        return {}
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return {}
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        summary_data = {}
        for row in ws.iter_rows(values_only=True):
            if row and len(row) >= 2:
                label = str(row[0]).strip().lower() if row[0] else ''
                value = row[1]
                if 'attachment' in label:
                    try:
                        summary_data['total_attachments'] = int(value) if value is not None else 0
                    except:
                        pass
                elif 'unique' in label or 'sender' in label:
                    try:
                        summary_data['unique_senders'] = int(value) if value is not None else 0
                    except:
                        pass
                elif 'total' in label and 'attachment' not in label:
                    try:
                        summary_data['total_count'] = int(value) if value is not None else 0
                    except:
                        pass
        wb.close()
        return summary_data
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_email_formatting__fba523fac66398bfd72f79ad1e012e07(env, config: Dict[str, Any]) -> Dict[str, bool]:
    """
    Extract email addresses and check if they are bolded.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key pointing to the docx file

    Returns:
        Dict mapping email addresses to boolean indicating if they are bold
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return {}
    with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        email_pattern = re.compile('\\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\\.[A-Z|a-z]{2,}\\b')
        email_formatting = {}
        for para in doc.paragraphs:
            for run in para.runs:
                emails = email_pattern.findall(run.text)
                for email in emails:
                    email_formatting[email] = run.font.bold is True
        return email_formatting
    finally:
        os.unlink(tmp_path)

def get_thunderbird_and_gdrive_pdfs__69a89619e7273cf1d986af62cad42166(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Composite getter that retrieves both Thunderbird email attachments and Google Drive files.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with keys:
            - email_subject: Subject of email to find in Thunderbird
            - settings_file: Path to Google Drive settings file
            - query: Query string to find files on Google Drive

    Returns:
        Dict with keys:
            - email_data: Data from Thunderbird email
            - gdrive_data: Data from Google Drive
    """
    import sys
    import os
    current_dir = os.path.dirname(os.path.abspath(__file__))
    parent_dir = os.path.dirname(current_dir)
    sys.path.insert(0, parent_dir)
    try:
        from getters.thunderbird import get_thunderbird_email_pdf_attachments__69a89619e7273cf1d986af62cad42166
    finally:
        sys.path.pop(0)
    email_config = {'email_subject': config.get('email_subject', 'Paper Recommendation'), 'folder': 'INBOX'}
    email_data = get_thunderbird_email_pdf_attachments__69a89619e7273cf1d986af62cad42166(env, email_config)
    gdrive_config = {'settings_file': config.get('settings_file', 'evaluation_examples/settings/googledrive/settings.yml'), 'query': config.get('query', "title contains '.pdf' and trashed = false and 'root' in parents")}
    gdrive_data = get_gdrive_file_count__69a89619e7273cf1d986af62cad42166(env, gdrive_config)
    return {'email_data': email_data, 'gdrive_data': gdrive_data}

def get_thunderbird_subject_date__2284321c(env, config: dict):
    """Extract subject and date columns from Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        list: List of tuples (subject, date) for data rows
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        data = []
        for row_idx in range(2, ws.max_row + 1):
            subject = ws.cell(row_idx, 1).value
            date = ws.cell(row_idx, 2).value
            data.append((subject, date))
        return data
    finally:
        os.unlink(tmp_path)

def get_thunderbird_cc_emails__f218f3c2(env, config: dict):
    """Extract sender and subject from Excel AND verify against Thunderbird mailbox.

    Reads the 2-column Excel file (sender, subject) and independently verifies
    these emails match the CC-filtered subset of the 5 most recent emails from
    Thunderbird's daily folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: {
            'excel_data': List of tuples (sender_name, subject) from Excel,
            'thunderbird_cc_emails': List of tuples (sender, subject) from Thunderbird,
            'valid': bool indicating if Excel matches Thunderbird data
        }
    """
    excel_data = _read_excel_file(env, config['path'])
    thunderbird_cc_emails = _get_thunderbird_cc_emails_from_mailbox(env)
    valid = _verify_excel_matches_thunderbird(excel_data, thunderbird_cc_emails)
    return {'excel_data': excel_data, 'thunderbird_cc_emails': thunderbird_cc_emails, 'valid': valid}

def get_thunderbird_folder_emails__4bba48e3(env, config: Dict[str, Any]) -> List[Dict[str, str]]:
    """
    Get emails from a specific Thunderbird folder.

    This function reads the Thunderbird folder mbox file and extracts email information
    including subject lines. The folder is stored as an mbox file in the profile directory.

    Args:
        env: Environment object
        config: Configuration dict with:
            - profile_path: Path to Thunderbird profile (e.g., /home/user/.thunderbird/xxx.default-release)
            - folder: Folder name (e.g., 'Bills')

    Returns:
        List[Dict[str, str]]: List of email dictionaries with 'subject' and 'raw' keys
    """
    profile_path = config.get('profile_path')
    folder_name = config.get('folder', 'Bills')
    if not profile_path:
        logger.error('No profile_path provided in config')
        return []
    local_folders_path = os.path.join(profile_path, 'Mail', 'Local Folders')
    folder_file = os.path.join(local_folders_path, folder_name)
    logger.info(f'Reading Thunderbird folder: {folder_file}')
    if not os.path.exists(folder_file):
        logger.warning(f'Folder file does not exist: {folder_file}')
        return []
    try:
        with open(folder_file, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
    except Exception as e:
        logger.error(f'Error reading folder file {folder_file}: {e}')
        return []
    emails = []
    messages = re.split('\\n(?=From - )', content)
    for msg in messages:
        if not msg.strip():
            continue
        if not msg.startswith('From - '):
            if 'Subject:' not in msg:
                continue
        subject_match = re.search('^Subject:\\s*(.+?)$', msg, re.MULTILINE | re.IGNORECASE)
        subject = subject_match.group(1).strip() if subject_match else ''
        status_match = re.search('X-Mozilla-Status:\\s*([0-9A-Fa-f]+)', msg)
        if status_match:
            status_code = status_match.group(1)
            if status_code in ['0008', '0009']:
                logger.debug(f'Skipping deleted message with subject: {subject}')
                continue
        emails.append({'subject': subject, 'raw': msg})
    logger.info(f"Found {len(emails)} emails in folder '{folder_name}'")
    return emails

def get_thunderbird_domain_count__1744c07b(env, config: dict):
    """Extract domain and count from Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary mapping domains to email counts
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        domain_counts = {}
        for row_idx in range(2, ws.max_row + 1):
            domain = ws.cell(row_idx, 1).value
            count = ws.cell(row_idx, 2).value
            if domain and count is not None:
                domain_counts[domain] = count
        return domain_counts
    finally:
        os.unlink(tmp_path)

def get_thunderbird_oldest_three__0c7bac1c(env, config: dict):
    """Extract sender and subject from the oldest 3 emails, with Thunderbird verification.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        dict: Dictionary with 'headers', 'data', and 'thunderbird_emails' keys
              - headers: tuple of column headers (col1, col2)
              - data: list of tuples (sender_name, subject) from Excel
              - thunderbird_emails: list of tuples (sender, subject, date) from daily folder, sorted oldest first
    """
    file_bytes = env.controller.get_file(config['path'])
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        if ws.max_row < 1:
            excel_result = {'headers': None, 'data': []}
        else:
            header1 = ws.cell(1, 1).value
            header2 = ws.cell(1, 2).value
            headers = (header1, header2)
            data = []
            for row_idx in range(2, min(5, ws.max_row + 1)):
                sender = ws.cell(row_idx, 1).value
                subject = ws.cell(row_idx, 2).value
                if sender is not None or subject is not None:
                    data.append((sender, subject))
            excel_result = {'headers': headers, 'data': data}
    finally:
        os.unlink(tmp_path)
    thunderbird_emails = []
    try:
        profile_path = '/home/user/.thunderbird'
        daily_mbox_path = None
        mail_local_path = os.path.join(profile_path, 'gah8lzo3.default-release/Mail/Local Folders')
        if os.path.exists(mail_local_path):
            for filename in os.listdir(mail_local_path):
                if filename.lower() == 'daily' or filename.lower().startswith('daily'):
                    candidate_path = os.path.join(mail_local_path, filename)
                    if not filename.endswith('.msf') and os.path.isfile(candidate_path):
                        daily_mbox_path = candidate_path
                        break
        if daily_mbox_path and os.path.exists(daily_mbox_path):
            try:
                mbox_bytes = env.controller.get_file(daily_mbox_path)
                with tempfile.NamedTemporaryFile(delete=False, mode='wb') as mbox_tmp:
                    mbox_tmp.write(mbox_bytes)
                    mbox_tmp_path = mbox_tmp.name
                try:
                    mbox = mailbox.mbox(mbox_tmp_path)
                    for message in mbox:
                        sender = message.get('From', '')
                        if '<' in sender:
                            sender = sender.split('<')[0].strip()
                        subject = message.get('Subject', '')
                        date_str = message.get('Date', '')
                        try:
                            from email.utils import parsedate_to_datetime
                            date_obj = parsedate_to_datetime(date_str)
                        except:
                            date_obj = datetime(1970, 1, 1)
                        thunderbird_emails.append((sender, subject, date_obj))
                    thunderbird_emails.sort(key=lambda x: x[2])
                finally:
                    os.unlink(mbox_tmp_path)
            except Exception as e:
                pass
    except Exception as e:
        pass
    excel_result['thunderbird_emails'] = thunderbird_emails
    return excel_result

def get_thunderbird_txt_backup__dd834e2c641c360491704c9db4c6f362(env, config):
    """
    Get list of .txt files in the backup directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to backup directory

    Returns:
        str: Path to cache file containing directory listing
    """
    backup_path = config.get('path', '/home/user/emails.txt.bak')
    result = env.controller.run_bash_script(f'ls -R {backup_path}', timeout=30)
    if result['returncode'] != 0:
        return None
    cache_path = os.path.join(env.cache_dir, 'emails.txt.bak.ls')
    with open(cache_path, 'w') as f:
        f.write(result['output'])
    return cache_path

def get_thunderbird_prefs__616461c5ba008feaf990bbd287063304(env, config):
    """
    Get Thunderbird prefs.js file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to downloaded prefs.js file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.js', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_thunderbird_prefs__a3ecd62701cd6d8575ca51c05b0fdb89(env, config):
    """
    Get Thunderbird prefs.js file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to downloaded prefs.js file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.js', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_thunderbird_date_from_subject__a92c1922(env, config: dict):
    """Extract dates from subject lines in Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of date strings extracted from subjects
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        dates = []
        for row_idx in range(2, ws.max_row + 1):
            date_str = ws.cell(row_idx, 1).value
            if date_str:
                dates.append(date_str)
        return dates
    finally:
        os.unlink(tmp_path)

def get_thunderbird_prefs__51cc10d9d3791de4bc35e6d331e3ef94(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the Thunderbird prefs.js file from the VM.

    Config:
        prefs_path (str): absolute path to the prefs.js file on the VM

    Returns:
        str: Local path to the downloaded prefs.js file, or None if failed
    """
    prefs_path = config.get('prefs_path')
    if not prefs_path:
        logger.error('prefs_path not provided in config')
        return None
    try:
        file_bytes = env.controller.get_file(prefs_path)
        if file_bytes is None:
            logger.warning(f'Failed to get prefs.js from VM: {prefs_path}')
            return None
        cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
        os.makedirs(env.cache_dir, exist_ok=True)
        with open(cache_path, 'wb') as f:
            f.write(file_bytes)
        logger.info(f'Successfully saved prefs.js: {cache_path} ({len(file_bytes)} bytes)')
        return cache_path
    except Exception as e:
        logger.error(f'Error getting Thunderbird prefs.js: {e}')
        return None

def get_thunderbird_prefs__5b041e31e4566087bb3bbe0b347cd7ff(env, config: Dict[str, Any]) -> str:
    """
    Get Thunderbird preferences file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        Path to the downloaded prefs.js file in cache directory

    Raises:
        ValueError: If file download fails
    """
    path = config['path']
    dest = config.get('dest', 'thunder-prefs.js')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        raise ValueError(f'Failed to download Thunderbird prefs from {path}')
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_prefs__bd92e86076085c3fd0cc432ab151ca8f(env, config):
    """
    Get Thunderbird prefs.js file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to downloaded prefs.js file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.js', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_thunderbird_prefs__cd29a448(env, config):
    """Get Thunderbird prefs.js file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_tb_filtered_emails__01abb20597fbc4723613687c7c60c0d3(env, config: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Extract filtered email list from a LibreOffice Calc spreadsheet.
    Expected to contain emails filtered by some criteria (e.g., keyword in subject).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        List of email dictionaries
    """
    file_path = config.get('path', '')
    if not file_path:
        return []
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return []
    with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        emails = []
        for (row_idx, row) in enumerate(ws.iter_rows(min_row=1, values_only=True), start=1):
            if row_idx == 1 and row and isinstance(row[0], str):
                if any((keyword in str(row[0]).lower() for keyword in ['sender', 'subject', 'from'])):
                    continue
            if row and any((cell is not None and str(cell).strip() for cell in row)):
                email_data = {'sender_name': str(row[0]).strip() if row[0] is not None else '', 'sender_email': str(row[1]).strip() if len(row) > 1 and row[1] is not None else '', 'subject': str(row[2]).strip() if len(row) > 2 and row[2] is not None else ''}
                emails.append(email_data)
        wb.close()
        return emails
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)

def get_default_email_client__b929c586(env, config: dict):
    """Gets the default email client application.

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        The default email client registered for x-scheme-handler/mailto
    """
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'x-scheme-handler/mailto']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_email_count_file__c95564e31fba18e60360c502646bdd5f(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Read the content of a text file that should contain the email count.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' (file path to read)

    Returns:
        Content of the file as string, or None if file doesn't exist
    """
    file_path = config.get('path', '/home/user/email_count.txt')
    file_bytes = env.controller.get_file(file_path)
    if file_bytes is None:
        return None
    try:
        content = file_bytes.decode('utf-8').strip()
        return content
    except Exception:
        return None

def get_thunderbird_folder_and_email_count__a91026e2(env, config: dict):
    """Check if folder exists and count emails in Bills folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - profile_path: Path to Thunderbird profile
            - folder_to_check: Name of folder to check existence
            - bills_folder_path: Path to Bills folder file

    Returns:
        dict: {
            "folder_exists": bool,
            "bills_email_count": int
        }
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    folder_to_check = config.get('folder_to_check', 'Archive')
    bills_folder_path = config.get('bills_folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills')
    result = {'folder_exists': False, 'bills_email_count': 0}
    folder_path = os.path.join(profile_path, 'Mail', 'Local Folders', folder_to_check)
    folder_file_bytes = env.controller.get_file(folder_path)
    if folder_file_bytes is not None:
        result['folder_exists'] = True
        logger.info(f'Folder {folder_to_check} exists')
    else:
        logger.info(f'Folder {folder_to_check} does not exist')
    bills_bytes = env.controller.get_file(bills_folder_path)
    if bills_bytes is not None:
        bills_content = bills_bytes.decode('utf-8', errors='ignore')
        emails = bills_content.split('FROM - ')
        for email in emails:
            if email.strip() and (not re.search('X-Mozilla-Status: 000[89]', email)):
                result['bills_email_count'] += 1
        logger.info(f"Found {result['bills_email_count']} non-deleted emails in Bills folder")
    else:
        logger.info('Bills folder not found')
    return result

def get_thunderbird_prefs__17405a90(env, config):
    """Get Thunderbird prefs.js file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_sender_count__44f17dcf(env, config: dict):
    """Extract sender name and email count from both Thunderbird mailbox and Excel report.

    This getter validates both:
    1. The actual Thunderbird mailbox data (5 most recent emails from daily folder)
    2. The Excel report created by the user

    This ensures the user actually analyzed Thunderbird emails rather than manually
    creating an Excel file with the expected values.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key for Excel file

    Returns:
        dict: {
            'excel_sender_counts': dict mapping sender names to counts from Excel,
            'thunderbird_sender_counts': dict mapping sender names to counts from Thunderbird,
            'excel_valid': bool indicating if Excel file is valid,
            'thunderbird_valid': bool indicating if Thunderbird data was read successfully
        }
    """
    result = {'excel_sender_counts': {}, 'thunderbird_sender_counts': {}, 'excel_valid': False, 'thunderbird_valid': False}
    try:
        file_bytes = env.controller.get_file(config['path'])
        import tempfile
        import os
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        try:
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            ws = wb.active
            if ws.max_column == 2:
                header1 = ws.cell(1, 1).value
                header2 = ws.cell(1, 2).value
                if header1 and header2 and isinstance(header1, str) and isinstance(header2, str):
                    h1_lower = header1.lower().strip()
                    h2_lower = header2.lower().strip()
                    sender_keywords = ['sender', 'from', 'email', 'name', 'address']
                    count_keywords = ['count', 'number', 'total', 'quantity', 'amount']
                    h1_is_sender = any((kw in h1_lower for kw in sender_keywords))
                    h1_is_count = any((kw in h1_lower for kw in count_keywords))
                    h2_is_sender = any((kw in h2_lower for kw in sender_keywords))
                    h2_is_count = any((kw in h2_lower for kw in count_keywords))
                    if h1_is_sender and h2_is_count or (h1_is_count and h2_is_sender):
                        sender_col = 1 if h1_is_sender else 2
                        count_col = 2 if h1_is_sender else 1
                        excel_counts = {}
                        for row_idx in range(2, ws.max_row + 1):
                            sender = ws.cell(row_idx, sender_col).value
                            count = ws.cell(row_idx, count_col).value
                            if sender and count is not None:
                                try:
                                    count_int = int(count)
                                    excel_counts[str(sender).strip()] = count_int
                                except (ValueError, TypeError):
                                    continue
                        if excel_counts:
                            result['excel_sender_counts'] = excel_counts
                            result['excel_valid'] = True
        finally:
            os.unlink(tmp_path)
    except Exception:
        pass
    try:
        thunderbird_data = get_thunderbird_recent_emails(env, top_n=5)
        if thunderbird_data and thunderbird_data.get('sender_counts'):
            result['thunderbird_sender_counts'] = thunderbird_data['sender_counts']
            result['thunderbird_valid'] = True
    except Exception:
        pass
    return result

def get_thunderbird_prefs__4d8c063289348a4e74b61401db33c326(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get Thunderbird preferences file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' and 'dest' keys

    Returns:
        Path to the downloaded prefs.js file in cache directory, or None if download fails
    """
    path = config['path']
    dest = config.get('dest', 'thunder-prefs.js')
    file_bytes = env.controller.get_file(path)
    if not file_bytes:
        return None
    cache_path = os.path.join(env.cache_dir, dest)
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_folder_counts__4f3b75dbef7eff0e54cd50b24ff5fb5e(env, config):
    """Get email metadata from multiple Thunderbird folders.

    Args:
        env: DesktopEnv instance
        config: Dict with 'paths' list of folder file paths

    Returns:
        dict: Mapping of folder path to list of email metadata dicts
              Each email dict contains: subject, date, sender, date_timestamp
    """
    paths = config.get('paths', [])
    if not paths:
        logger.error('No paths specified in config')
        return None
    folder_emails = {}
    for path in paths:
        try:
            file_bytes = env.controller.get_file(path)
            if not file_bytes:
                logger.info(f'File not found: {path}')
                folder_emails[path] = []
                continue
            content = file_bytes.decode('utf-8', errors='ignore')
            messages = content.split('From - ')
            emails = []
            for msg in messages[1:]:
                if not msg.strip():
                    continue
                if re.search('X-Mozilla-Status: 000[89]', msg):
                    continue
                email_data = {}
                subject_match = re.search('^Subject:\\s*(.*)$', msg, re.MULTILINE)
                if subject_match:
                    email_data['subject'] = subject_match.group(1).strip()
                else:
                    email_data['subject'] = ''
                from_match = re.search('^From:\\s*(.*)$', msg, re.MULTILINE)
                if from_match:
                    email_data['sender'] = from_match.group(1).strip()
                else:
                    email_data['sender'] = ''
                date_match = re.search('^Date:\\s*(.*)$', msg, re.MULTILINE)
                if date_match:
                    date_str = date_match.group(1).strip()
                    email_data['date'] = date_str
                    try:
                        dt = parsedate_to_datetime(date_str)
                        email_data['date_timestamp'] = dt.timestamp()
                    except Exception as e:
                        logger.warning(f"Could not parse date '{date_str}': {e}")
                        email_data['date_timestamp'] = 0
                else:
                    email_data['date'] = ''
                    email_data['date_timestamp'] = 0
                emails.append(email_data)
            folder_emails[path] = emails
            logger.info(f'Found {len(emails)} messages in {path}')
        except Exception as e:
            logger.error(f'Error reading {path}: {e}')
            folder_emails[path] = None
    return folder_emails

def get_txt_contains_email__a3b9c6d3(env, config):
    """
    Check if text file contains all required contact information.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path', 'email', and 'twitter' keys

    Returns:
        bool: True if all contact information is found, False otherwise
    """
    import re
    file_path = config.get('path', '')
    expected_email = config.get('email', '')
    expected_twitter = config.get('twitter', '')
    try:
        file_bytes = env.controller.get_file(file_path)
        if file_bytes is None:
            return False
        content = file_bytes.decode('utf-8', errors='ignore').lower()
        email_found = expected_email.lower() in content
        twitter_found = expected_twitter.lower() in content
        return email_found and twitter_found
    except Exception as e:
        return False

def get_thunderbird_and_chrome_state__0f025d55a86976e1ae415a210c74f719(env, config: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get comprehensive state from both Thunderbird and Chrome to verify the complete workflow.

    This getter collects:
    1. Thunderbird accessibility tree state (to verify Bills folder was accessed)
    2. Chrome history with timestamps (to verify amazon.com was recently visited)
    3. Chrome open tabs (to verify tab count and URLs)

    Args:
        env: Environment object
        config: Configuration dict

    Returns:
        dict: {
            'thunderbird_state': {
                'bills_folder_visible': bool,
                'bills_folder_selected': bool,
                'email_displayed': bool,
                'folder_name': str or None
            },
            'chrome_history': [
                {'url': str, 'title': str, 'last_visit_time': int},
                ...
            ],
            'chrome_tabs': [
                {'url': str, 'title': str},
                ...
            ]
        }
    """
    logger.info('[THUNDERBIRD_CHROME_STATE] Starting comprehensive state collection')
    result = {'thunderbird_state': {}, 'chrome_history': [], 'chrome_tabs': []}
    try:
        logger.info('[THUNDERBIRD_CHROME_STATE] Collecting Thunderbird accessibility tree state')
        thunderbird_state = _get_thunderbird_accessibility_state(env)
        result['thunderbird_state'] = thunderbird_state
        logger.info(f'[THUNDERBIRD_CHROME_STATE] Thunderbird state: {thunderbird_state}')
    except Exception as e:
        logger.error(f'[THUNDERBIRD_CHROME_STATE] Error getting Thunderbird state: {e}')
        result['thunderbird_state'] = {'bills_folder_visible': False, 'email_displayed': False, 'folder_name': None, 'error': str(e)}
    try:
        logger.info('[THUNDERBIRD_CHROME_STATE] Collecting Chrome history')
        chrome_history = _get_chrome_history_recent(env, limit=50)
        result['chrome_history'] = chrome_history
        logger.info(f'[THUNDERBIRD_CHROME_STATE] Found {len(chrome_history)} recent history entries')
    except Exception as e:
        logger.error(f'[THUNDERBIRD_CHROME_STATE] Error getting Chrome history: {e}')
        result['chrome_history'] = []
    try:
        logger.info('[THUNDERBIRD_CHROME_STATE] Collecting Chrome open tabs')
        from desktop_env.evaluators.getters.chrome import get_open_tabs_info
        chrome_tabs = get_open_tabs_info(env, {})
        result['chrome_tabs'] = chrome_tabs if chrome_tabs else []
        logger.info(f"[THUNDERBIRD_CHROME_STATE] Found {len(result['chrome_tabs'])} open tabs")
    except Exception as e:
        logger.error(f'[THUNDERBIRD_CHROME_STATE] Error getting Chrome tabs: {e}')
        result['chrome_tabs'] = []
    logger.info('[THUNDERBIRD_CHROME_STATE] State collection complete')
    return result

def get_email_attachments__d38192b0(env, config: dict):
    """Extract attachments from a Thunderbird email draft.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with parameters including:
            - subject_title: The email subject to identify the draft
            - expected_attachment: The expected attachment filename

    Returns:
        str: Output from the attachment verification script
    """
    subject_title = config.get('subject_title', '')
    expected_attachment = config.get('expected_attachment', '')
    result = env.controller.run_bash_script(f"python /home/user/show-thunderbird-attachments.py '{subject_title}' '{expected_attachment}'", timeout=30)
    if result['returncode'] == 0:
        return result['output'].strip()
    else:
        return None

def get_test_emails_backup__a9ff16fc1be9331aa02283bf7c168b3e(env, config):
    """
    Get list of files in the test emails backup directory.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to backup directory

    Returns:
        str: Path to cache file containing directory listing
    """
    backup_path = config.get('path', '/home/user/test_emails.bak')
    result = env.controller.run_bash_script(f'ls -R {backup_path}', timeout=30)
    if result['returncode'] != 0:
        return None
    cache_path = os.path.join(env.cache_dir, 'test_emails.bak.ls')
    with open(cache_path, 'w') as f:
        f.write(result['output'])
    return cache_path

def get_thunderbird_prefs__550830d3(env, config):
    """Get Thunderbird prefs.js file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_prefs__e06bf170da6e0042eda0ceb7f1f9f833(env, config):
    """
    Get Thunderbird prefs.js file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to downloaded prefs.js file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.js', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_thunderbird_prefs__6e76242e9ee0066f62efb661397a7719(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the path to the Thunderbird preferences file (prefs.js).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
                - 'prefs_path': Path to the prefs.js file

    Returns:
        str: Path to the prefs.js file, or None if not found
    """
    prefs_path = config.get('prefs_path')
    if not prefs_path:
        logger.error('No prefs_path specified in config')
        return None
    check_cmd = f"test -f {prefs_path} && echo 'exists' || echo 'not found'"
    result = env.controller.run_bash_script(check_cmd, timeout=10)
    if result.get('returncode') != 0 or 'exists' not in result.get('output', ''):
        logger.error(f'Thunderbird prefs file not found at {prefs_path}')
        return None
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        logger.error(f'Failed to retrieve prefs file from {prefs_path}')
        return None
    import tempfile
    import os
    (fd, temp_path) = tempfile.mkstemp(suffix='_prefs.js', prefix='thunderbird_')
    try:
        os.write(fd, file_bytes)
    finally:
        os.close(fd)
    logger.info(f'Retrieved Thunderbird prefs from {prefs_path} to {temp_path}')
    return temp_path

def get_thunderbird_prefs__dddfffe093b014a09b5e9ec4efaf461f(env, config: Dict[str, Any]) -> Optional[str]:
    """
    Get the path to the Thunderbird preferences file (prefs.js).

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
                - 'prefs_path': Path to the prefs.js file

    Returns:
        str: Path to the prefs.js file, or None if not found
    """
    prefs_path = config.get('prefs_path')
    if not prefs_path:
        logger.error('No prefs_path specified in config')
        return None
    check_cmd = f"test -f {prefs_path} && echo 'exists' || echo 'not found'"
    result = env.controller.run_bash_script(check_cmd, timeout=10)
    if result.get('returncode') != 0 or 'exists' not in result.get('output', ''):
        logger.error(f'Thunderbird prefs file not found at {prefs_path}')
        return None
    file_bytes = env.controller.get_file(prefs_path)
    if not file_bytes:
        logger.error(f'Failed to retrieve prefs file from {prefs_path}')
        return None
    import tempfile
    import os
    (fd, temp_path) = tempfile.mkstemp(suffix='_prefs.js', prefix='thunderbird_')
    try:
        os.write(fd, file_bytes)
    finally:
        os.close(fd)
    logger.info(f'Retrieved Thunderbird prefs from {prefs_path} to {temp_path}')
    return temp_path

def get_thunderbird_email_read_status__a63fb94e(env, config: dict):
    """Check read status of all emails in Bills folder.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - bills_folder_path: Path to Bills folder file

    Returns:
        dict: {
            "all_read": bool,
            "total_emails": int,
            "read_emails": int
        }
    """
    bills_folder_path = config.get('bills_folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills')
    result = {'all_read': False, 'total_emails': 0, 'read_emails': 0}
    bills_bytes = env.controller.get_file(bills_folder_path)
    if bills_bytes is not None:
        bills_content = bills_bytes.decode('utf-8', errors='ignore')
        emails = bills_content.split('FROM - ')
        for email in emails:
            if email.strip() and (not re.search('X-Mozilla-Status: 000[89]', email)):
                result['total_emails'] += 1
                if re.search('X-Mozilla-Status: 000[13]', email):
                    result['read_emails'] += 1
        if result['total_emails'] > 0 and result['read_emails'] == result['total_emails']:
            result['all_read'] = True
        logger.info(f"Total emails: {result['total_emails']}, Read emails: {result['read_emails']}")
    else:
        logger.info('Bills folder not found')
    return result

def get_email_subjects_file__df60c313c1f3b2e712e9756c47386ddc(env, config):
    """
    Get the content of email subjects list file.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' to subjects file

    Returns:
        str: Content of the subjects file
    """
    file_path = config.get('path', '/home/user/email_subjects.txt')
    file_bytes = env.controller.get_file(file_path)
    if not file_bytes:
        return None
    content = file_bytes.decode('utf-8', errors='ignore')
    return content

def get_thunderbird_subject__c9ce3f52(env, config):
    """Get the email subject from Thunderbird compose window.

    Args:
        env: DesktopEnv instance
        config: Configuration dict (unused but required by framework)

    Returns:
        str: The subject text or empty string if not found
    """
    accessibility_tree = env.controller.get_accessibility_tree()
    if not accessibility_tree:
        return ''
    lines = accessibility_tree.split('\n')
    for (i, line) in enumerate(lines):
        if 'name=Subject' in line and 'entry' in line.lower():
            if 'text=' in line:
                parts = line.split('text=')
                if len(parts) > 1:
                    value_part = parts[1]
                    if '"' in value_part:
                        return value_part.split('"')[1]
            for j in range(i, min(i + 5, len(lines))):
                if 'text=' in lines[j]:
                    parts = lines[j].split('text=')
                    if len(parts) > 1:
                        value_part = parts[1]
                        if '"' in value_part:
                            return value_part.split('"')[1]
    return ''

def get_default_email_client__119177bf(env, config: dict):
    """Gets the default email client on Linux."""
    os_type = env.vm_platform
    if os_type == 'Linux':
        app = get_vm_command_line(env, {'command': ['xdg-mime', 'query', 'default', 'x-scheme-handler/mailto']})
        if app:
            return app
        else:
            return 'unknown'
    elif os_type == 'Darwin':
        raise Exception('Unsupported operating system', os_type)
    elif os_type == 'Windows':
        raise Exception('Unsupported operating system', os_type)
    else:
        raise Exception('Unsupported operating system', os_type)

def get_thunderbird_email_count_and_folder__3c15b6d0(env, config: dict):
    """Count emails in Bills folder and check if target folder exists.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with:
            - profile_path: Path to Thunderbird profile
            - bills_folder_path: Path to Bills folder file
            - folder_to_check: Name of folder to check existence

    Returns:
        dict: {
            "email_count": int,
            "folder_exists": bool
        }
    """
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    bills_folder_path = config.get('bills_folder_path', '/home/user/.thunderbird/t5q2a5hp.default-release/Mail/Local Folders/Bills')
    folder_to_check = config.get('folder_to_check', 'Processed')
    result = {'email_count': 0, 'folder_exists': False}
    bills_bytes = env.controller.get_file(bills_folder_path)
    if bills_bytes is not None:
        bills_content = bills_bytes.decode('utf-8', errors='ignore')
        emails = bills_content.split('FROM - ')
        for email in emails:
            if email.strip() and (not re.search('X-Mozilla-Status: 000[89]', email)):
                result['email_count'] += 1
        logger.info(f"Found {result['email_count']} non-deleted emails in Bills folder")
    else:
        logger.info('Bills folder not found')
    folder_path = os.path.join(profile_path, 'Mail', 'Local Folders', folder_to_check)
    folder_file_bytes = env.controller.get_file(folder_path)
    if folder_file_bytes is not None:
        result['folder_exists'] = True
        logger.info(f'Folder {folder_to_check} exists')
    else:
        logger.info(f'Folder {folder_to_check} does not exist')
    return result

def get_thunderbird_filters__f57034f8(env, config: dict):
    """Extract Thunderbird message filter definitions.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'filter_file_path'

    Returns:
        str: Path to downloaded filter file
    """
    filter_path = config.get('filter_file_path', '/home/user/.thunderbird/t5q2a5hp.default-release/ImapMail/outlook.office365.com/msgFilterRules.dat')
    result = env.controller.get_file(filter_path)
    if result is None:
        return None
    cache_path = os.path.join(env.cache_dir, 'msgFilterRules.dat')
    with open(cache_path, 'wb') as f:
        f.write(result)
    return cache_path

def get_thunderbird_filter_rules__8a2db6ad(env, config):
    """Get Thunderbird msgFilterRules.dat file from the VM."""
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    filter_path = f'{profile_path}/ImapMail/outlook.office365.com/msgFilterRules.dat'
    file_bytes = env.controller.get_file(filter_path)
    if file_bytes is None:
        return None
    import os
    cache_path = os.path.join(env.cache_dir, 'msgFilterRules.dat')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    return cache_path

def get_thunderbird_three_cols__07170c7f(env, config: dict):
    """Extract sender, subject, and has_cc from Excel report.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        list: List of tuples (sender_name, subject, has_cc)
    """
    file_bytes = env.controller.get_file(config['path'])
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        ws = wb.active
        data = []
        for row_idx in range(2, ws.max_row + 1):
            sender = ws.cell(row_idx, 1).value
            subject = ws.cell(row_idx, 2).value
            has_cc = ws.cell(row_idx, 3).value
            data.append((sender, subject, has_cc))
        return data
    finally:
        os.unlink(tmp_path)

def get_thunderbird_prefs__734b4908962171d93d8cfb6878a651d2(env, config):
    """
    Get Thunderbird prefs.js file from VM.

    Args:
        env: DesktopEnv instance
        config: Configuration dict with 'path' key

    Returns:
        str: Path to downloaded prefs.js file
    """
    file_bytes = env.controller.get_file(config['path'])
    if not file_bytes:
        return None
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='wb', suffix='.js', delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    return tmp_path

def get_thunderbird_prefs__497f2c49(env, config):
    """
    Get Thunderbird prefs.js file from the VM and extract BCC settings for the specified email account.

    Returns a dict with:
    - identity_id: The identity ID for the email account
    - do_bcc: Boolean indicating if BCC is enabled
    - bcc_list: The BCC address list (if configured)
    - email: The email account being checked
    """
    import os
    import json
    profile_path = config.get('profile_path', '/home/user/.thunderbird/t5q2a5hp.default-release')
    prefs_path = f'{profile_path}/prefs.js'
    file_bytes = env.controller.get_file(prefs_path)
    if file_bytes is None:
        return None
    cache_path = os.path.join(env.cache_dir, 'thunderbird_prefs.js')
    with open(cache_path, 'wb') as f:
        f.write(file_bytes)
    result = {'identity_id': None, 'do_bcc': None, 'bcc_list': None, 'email': 'anonym-x2024@outlook.com', 'prefs_path': cache_path}
    try:
        with open(cache_path, 'r') as f:
            content = f.read()
        for line in content.splitlines():
            if 'mail.identity.id' in line and '.useremail' in line and ('user_pref' in line):
                if 'anonym-x2024@outlook.com' in line:
                    start = line.find('mail.identity.id') + len('mail.identity.id')
                    end = line.find('.useremail')
                    identity_id = line[start:end]
                    result['identity_id'] = identity_id
                    break
        if result['identity_id']:
            identity_id = result['identity_id']
            for line in content.splitlines():
                if f'mail.identity.id{identity_id}.doBcc' in line and 'user_pref' in line:
                    parts = line.split(', ')
                    if len(parts) >= 2:
                        value_part = parts[1].strip().rstrip(');')
                        result['do_bcc'] = json.loads(value_part)
                    break
            for line in content.splitlines():
                if f'mail.identity.id{identity_id}.doBccList' in line and 'user_pref' in line:
                    parts = line.split(', ')
                    if len(parts) >= 2:
                        value_part = parts[1].strip().rstrip(');')
                        result['bcc_list'] = json.loads(value_part)
                    break
    except Exception as e:
        import logging
        logger = logging.getLogger('desktopenv.getter.thunderbird')
        logger.error(f'Error parsing prefs.js: {e}')
        return None
    return result
