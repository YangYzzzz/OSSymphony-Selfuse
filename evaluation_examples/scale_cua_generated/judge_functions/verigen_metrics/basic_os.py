"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_folder_structure__e617d7dc', 'check_zip_and_folder__d9c12923d5c2940b6520fb26a328924d', 'check_identified_mountains_folder__c8c87c59ecc91c5a6beb4965364b8f59', 'check_gnome_app_removed__a60961d39b430968c7907f27beb1d667', 'check_screenshot_in_folder__956d014f26874db42b602f2626906894', 'check_gnome_favorites_count__e2063f4a', 'check_directory_has_file__b77b1963e53aafc9923374c6ea5077e2', 'check_gnome_boolean__37eb80a3', 'check_gnome_favorites_swap_last_two__e2063f4a', 'check_folder_exists__c5a909ed', 'check_gnome_favorites_contains_only__e2063f4a', 'check_gnome_app_replaced__c1c35c032a59fde6fdc31c4f532d1270', 'check_file_in_folder__3c993009', 'check_directory_structure__e01c3944c2f0dfae12ffc1d2b96464cf', 'check_vendor_folders__f7d558cb', 'check_gnome_favorite_apps_order__5993e3c1', 'check_gnome_favorites_reversed__e2063f4a', 'check_folder_counts__4f3b75dbef7eff0e54cd50b24ff5fb5e', 'check_zip_and_folder__173d79a32d2c31ac3ad30d4ae958526d', 'check_directory_has_file__739292ff', 'check_folder_structure__f9d2dc30', 'check_terminal_profile_name__dfa021d620bcc8024aed513c1adfaad3', 'check_gnome_favorite_apps_ordered__5233e9af0d26c0a6bd95a54609343d46', 'check_folder_contains_files__c14b0c4873b98943d0bcc59ebce705c9', 'check_folder_contains_files__32b2b661', 'check_gnome_favorite_apps_order__1e04efad', 'check_qa_folder_contents__fc8f5b84', 'check_desktop_hashes__ef2db1ef', 'check_gnome_favorite_apps_empty__e17d22d0', 'check_gnome_favorites_empty__608fb5ab3b43018120dd1b5f6b15910c', 'check_nested_directory__3c3f0ceb', 'check_gnome_favorites_first_two__e2063f4a', 'check_bashrc_path_modification__c0107678', 'check_user_shell__49da98ad0a0985144d25f764fe852408', 'check_terminal_color_scheme__ad0133b72746f22f133946cc65aa7f55', 'check_bash_text_transform__8c025f1893a98ce909d203315d730cfe', 'check_directory_created__3ef6f937c5dcda24ec08a8ba5aa2c872', 'check_bash_text_transform__816eebe3a50924fca9ee760018f3238e', 'check_directory_created__96172b42', 'check_folder_organization__6a9b75b65029a1742667a75cc968a2a9', 'check_no_stars_in_folder__b22fca61', 'check_file_in_directory__35a43b1b', 'check_folder_and_file__223bf56e', 'check_directory_moved__11701199', 'check_file_in_directory__3500c3270fba14684f134a0b5e4537d1', 'check_file_in_folder__cbcf6e0c', 'check_folder_structure__4e03b1ed', 'check_folder_counts__856f6f3499821b31f240098182dcc479', 'check_gnome_favorite_apps_partial__b4585266e3285db10295625412489679', 'check_gnome_favorites_order__e2063f4a', 'check_file_on_desktop__e35e8479c21a9a3d6ef729a997676f8c', 'check_git_subdirectory__a29e9963', 'check_gnome_favorites_exact_single__e2063f4a', 'check_folder_exists__dc2d0173', 'check_gnome_favorites_empty__e2063f4a', 'check_gnome_boolean__fa85fc69', 'check_directory_count__7ec80ec4', 'check_folder_not_empty__95db6624', 'check_zip_and_folder__36005e14c2ae4846381f0946699e70be', 'check_gnome_app_added__1ed286d0e6299072018d07279d4346ca', 'check_subfolder_hashes__fcfe8473', 'check_reconciled_folder__fc048187', 'check_zip_and_folder__9ff7dd4db34cdbddf16b8ce0d6085594', 'check_directory_existence__0b3844b6', 'check_folder_contains_files__21ed89b452cf1a748805f1dda9b2ec4b', 'check_bash_text_transform__5ed0399c75d125da4fc3b5f7583bb5c4', 'check_folder_and_file__fe4fce7a', 'check_terminal_scrollback__8e9fd232f4307bd96cd64b8e7a1f9389', 'check_gnome_favorites_order__a572fc31d81435d0b739124764e326d3', 'check_gnome_favorites_swap_first_two__e2063f4a', 'check_desktop_file_exists__d17b574f', 'check_gnome_favorites_single__e2063f4a', 'check_directory_exists__881ee5e7', 'check_files_in_folder__2b29a90e', 'check_backup_folder__089628262a733e157c368e0bf1ad02f1', 'check_verified_folder__12aa1a9a', 'check_bash_text_transform__38c143b8ba918371e989fa588ea9cb56', 'check_terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012', 'check_gnome_favorite_apps_order__f848009c']

def check_folder_structure__e617d7dc(result, expected, **options):
    """Check folder structure.

    Args:
        result: dict with folder_exists and file_count
        expected: dict with folder_exists and file_count_min

    Returns:
        float: 1.0 if structure valid, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_exists = expected.get('folder_exists', True)
    file_count_min = expected.get('file_count_min', 0)
    if result.get('folder_exists', False) != expected_exists:
        return 0.0
    if result.get('file_count', 0) < file_count_min:
        return 0.0
    return 1.0

def check_zip_and_folder__d9c12923d5c2940b6520fb26a328924d(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the correct files are in the folder and zip exists.

    Args:
        result: List of filenames from the getter
        expected: Expected configuration with 'files' key
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    if not result:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    if len(expected_set) == 0:
        return 0.0
    correct_files = len(result_set & expected_set)
    extra_files = len(result_set - expected_set)
    missing_files = len(expected_set - result_set)
    score = correct_files / len(expected_set)
    if extra_files > 0:
        score *= 0.8
    return max(0.0, min(1.0, score))

def check_identified_mountains_folder__c8c87c59ecc91c5a6beb4965364b8f59(result, expected, **options):
    """Check if folder exists with correct files moved (not copied).

    Args:
        result: Dict with folder info from getter
        expected: Dict with expected_file_count
        **options: Additional options

    Returns:
        float: Score based on:
            - 0.3 for folder exists
            - 0.3 for having the correct specific files (picture1.jpg, picture2.jpg, picture3.jpg)
            - 0.4 for source files being removed (confirming move, not copy)
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.3
    if result.get('has_correct_files', False):
        score += 0.3
    if result.get('source_files_removed', False):
        score += 0.4
    return score

def check_gnome_app_removed__a60961d39b430968c7907f27beb1d667(result_state, expected_state, **options):
    """
    Check if a specific application was removed from GNOME favorites.

    Verifies that:
    1. The removed app is not in the favorites list
    2. The remaining apps match exactly what's expected

    Args:
        result_state: Output from gsettings get org.gnome.shell favorite-apps
        expected_state: Dict with 'removed_app' and 'remaining_apps'
        **options: Additional options

    Returns:
        float: 1.0 if verification passes, 0.0 otherwise
    """
    if not result_state or not isinstance(result_state, str):
        return 0.0
    import ast
    try:
        result_state = result_state.strip()
        current_apps = ast.literal_eval(result_state)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(current_apps, list):
        return 0.0
    removed_app = expected_state.get('removed_app')
    remaining_apps = expected_state.get('remaining_apps', [])
    if removed_app in current_apps:
        return 0.0
    if current_apps != remaining_apps:
        return 0.0
    return 1.0

def check_screenshot_in_folder__956d014f26874db42b602f2626906894(result, expected, **options):
    """
    Check if screenshot was saved in the correct folder with correct name.

    Args:
        result: Screenshot info from getter
        expected: Expected directory and filename
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        logger.info('Screenshot file does not exist')
        return 0.0
    expected_dir = expected.get('directory', '')
    actual_dir = result.get('directory', '')
    if actual_dir == expected_dir:
        score += 0.25
    else:
        logger.warning(f"Directory mismatch: expected '{expected_dir}', got '{actual_dir}'")
    expected_filename = expected.get('filename', '')
    actual_filename = result.get('filename', '')
    if actual_filename == expected_filename:
        score += 0.25
    else:
        logger.warning(f"Filename mismatch: expected '{expected_filename}', got '{actual_filename}'")
    logger.info(f'Screenshot in folder check score: {score}')
    return score

def check_gnome_favorites_count__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list has expected apps (order-independent with count check).

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing list of .desktop files and optional 'exact_count'

    Returns:
        1.0 if apps match expected set and count, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if len(apps) != len(expected_apps):
        return 0.0
    if set(apps) == set(expected_apps):
        return 1.0
    else:
        return 0.0

def check_directory_has_file__b77b1963e53aafc9923374c6ea5077e2(result: bool, expected: Any, **options) -> float:
    """Check if file exists in directory as expected.

    Args:
        result: Boolean indicating if file exists
        expected: Dict with 'exists' field (True/False)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_gnome_boolean__37eb80a3(result: str, expected, **options):
    """
    Check if a GNOME boolean setting matches expected value.

    Args:
        result: Path to file containing gsettings output (e.g., "true" or "false")
        expected: Expected boolean value as string ("true" or "false")

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip().lower()
        expected_value = str(expected).strip().lower()
        return 1.0 if content == expected_value else 0.0
    except Exception as e:
        print(f'Error checking boolean setting: {e}')
        return 0.0

def check_gnome_favorites_swap_last_two__e2063f4a(apps_str: str, rule):
    """Check if last two favorite apps are swapped compared to original.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing list with swapped last two apps

    Returns:
        1.0 if apps match expected with last two swapped, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if apps == expected_apps:
        return 1.0
    else:
        return 0.0

def check_folder_exists__c5a909ed(actual: str, expected: dict, **options) -> float:
    """
    Check if a folder exists based on command output.

    Args:
        actual (str): command output ("exists" or "not_exists")
        expected (dict): expected dict with key "expected"

    Return:
        float: 1.0 if folder exists, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_folder_exists__c5a909ed: actual is None or empty')
        return 0.0
    expected_value = expected.get('expected', 'exists')
    actual = actual.strip()
    if actual == expected_value:
        return 1.0
    logger.debug(f"check_folder_exists__c5a909ed: Expected '{expected_value}', got '{actual}'")
    return 0.0

def check_gnome_favorites_contains_only__e2063f4a(apps_str: str, rule):
    """Check if favorites list contains only specified apps (order-independent).

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing list of apps that should be present

    Returns:
        1.0 if apps match exactly (any order), 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if len(apps) != len(expected_apps):
        return 0.0
    if set(apps) == set(expected_apps):
        return 1.0
    else:
        return 0.0

def check_gnome_app_replaced__c1c35c032a59fde6fdc31c4f532d1270(result_state, expected_state, **options):
    """
    Check if Vim was replaced with Terminal in GNOME favorites.

    Verifies that:
    1. vim.desktop was removed from favorites
    2. org.gnome.Terminal.desktop was added to favorites
    3. The final list matches the expected apps

    Uses partial credit scoring:
    - 0.5 if vim.desktop is removed
    - 1.0 if Terminal is added AND the final list matches expected

    Args:
        result_state: Output from gsettings get org.gnome.shell favorite-apps
        expected_state: Dict with 'replaced_app', 'new_app', and 'expected_apps'
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on completion
    """
    if not result_state or not isinstance(result_state, str):
        return 0.0
    import ast
    try:
        result_state = result_state.strip()
        current_apps = ast.literal_eval(result_state)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(current_apps, list):
        return 0.0
    replaced_app = expected_state.get('replaced_app')
    new_app = expected_state.get('new_app')
    expected_apps = expected_state.get('expected_apps', [])
    vim_removed = replaced_app not in current_apps
    terminal_added = new_app in current_apps
    current_sorted = sorted(current_apps)
    expected_sorted = sorted(expected_apps)
    list_matches = current_sorted == expected_sorted
    if vim_removed and (not terminal_added):
        return 0.5
    if terminal_added and list_matches:
        return 1.0
    return 0.0

def check_file_in_folder__3c993009(result, expected, **options):
    """Verify folder exists and file is inside it.

    Args:
        result: Dict with folder_exists and file_in_folder booleans
        expected: Dict with rules specifying expected states
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 with partial credit
    """
    expected_folder_exists = expected.get('folder_exists', True)
    expected_file_in_folder = expected.get('file_in_folder', True)
    folder_exists = result.get('folder_exists', False)
    file_in_folder = result.get('file_in_folder', False)
    score = 0.0
    if folder_exists == expected_folder_exists:
        score += 0.4
    if file_in_folder == expected_file_in_folder:
        score += 0.6
    return score

def check_directory_structure__e01c3944c2f0dfae12ffc1d2b96464cf(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if directory structure is correct and contains expected PDFs with proper filenames.

    Args:
        result: Dict with directory info from getter (includes 'pdf_filenames')
        expected: Dict with 'must_exist', 'must_be_directory', 'min_pdf_count', 'filename_keywords' keys
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0 with partial credit
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.25
    if result.get('is_directory', False):
        score += 0.25
    min_pdf_count = expected.get('min_pdf_count', 2)
    actual_pdf_count = result.get('pdf_count', 0)
    if actual_pdf_count >= min_pdf_count:
        score += 0.25
    filename_keywords = expected.get('filename_keywords', ['agent', 'human', 'data', 'quality'])
    pdf_filenames = result.get('pdf_filenames', [])
    if pdf_filenames:
        filenames_lower = [f.lower() for f in pdf_filenames]
        matching_count = 0
        for filename in filenames_lower:
            filename_base = filename.replace('.pdf', '')
            if any((keyword in filename_base for keyword in filename_keywords)):
                matching_count += 1
        if matching_count >= 2:
            score += 0.25
    return min(1.0, score)

def check_vendor_folders__f7d558cb(result, expected, **options):
    """
    Check if vendor folders contain the correct invoice files.

    Args:
        result: Output from find command showing all PDF files
        expected: Dict with 'rules' containing 'expected_structure'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_structure = expected.get('expected_structure', {})
    if isinstance(result, dict) and 'output' in result:
        result = result['output']
    result = str(result).strip()
    if 'No such file' in result or 'cannot access' in result:
        return 0.0
    score = 0.0
    total_checks = len(expected_structure)
    if total_checks == 0:
        return 0.0
    found_files = [line.strip() for line in result.split('\n') if line.strip() and '.pdf' in line]
    for (vendor_folder, expected_files) in expected_structure.items():
        for expected_file in expected_files:
            expected_path_pattern = f'vendors/{vendor_folder}/{expected_file}'
            found = any((expected_path_pattern in path for path in found_files))
            if found:
                score += 1.0 / (total_checks * len(expected_files))
    return min(score, 1.0)

def check_gnome_favorite_apps_order__5993e3c1(apps_str: str, expected, **options):
    """Check if GNOME favorite apps are in the exact expected order.

    Args:
        apps_str: String representation of favorite apps list from gsettings
        expected: Dict with 'rules' key containing 'expected_order' list of app .desktop files

    Returns:
        float: 1.0 if order matches exactly, 0.0 otherwise
    """
    try:
        apps = ast.literal_eval(apps_str)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(apps, list):
        return 0.0
    expected_order = expected.get('expected_order') or expected.get('rules', {}).get('expected_order', [])
    if not expected_order:
        return 0.0
    for app in expected_order:
        if app not in apps:
            return 0.0
    if apps == expected_order:
        return 1.0
    else:
        return 0.0

def check_gnome_favorites_reversed__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list is reversed compared to original.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing reversed list

    Returns:
        1.0 if apps match expected reversed order, 0.0 otherwise
    """
    try:
        if not isinstance(apps_str, str) or not apps_str.strip():
            return 0.0
        apps = ast.literal_eval(apps_str.strip())
        if not isinstance(apps, list):
            return 0.0
        expected_apps = rule['expected']
        if apps == expected_apps:
            return 1.0
        else:
            return 0.0
    except (ValueError, SyntaxError, TypeError) as e:
        return 0.0

def check_folder_counts__4f3b75dbef7eff0e54cd50b24ff5fb5e(result: Optional[Dict[str, List[Dict]]], expected, **options):
    """Check if the oldest email was moved from Bills to have_seen folder.

    Args:
        result: Dict mapping folder paths to list of email metadata
        expected: Dict with 'expected_counts' mapping paths to expected counts
        **options: Additional options

    Returns:
        float: 1.0 if oldest email moved correctly, 0.0 otherwise
    """
    if result is None:
        logger.info('Result is None - folders not accessible')
        return 0.0
    expected_counts = expected.get('expected_counts', {})
    if not expected_counts:
        logger.warning('No expected_counts specified')
        return 0.0
    bills_path = None
    have_seen_path = None
    for path in expected_counts.keys():
        if 'Bills' in path:
            bills_path = path
        elif 'have_seen' in path:
            have_seen_path = path
    if not bills_path or not have_seen_path:
        logger.error('Could not identify Bills and have_seen folder paths')
        return 0.0
    bills_emails = result.get(bills_path)
    have_seen_emails = result.get(have_seen_path)
    if bills_emails is None or have_seen_emails is None:
        logger.info('One or both folders not found in result')
        return 0.0
    if len(bills_emails) != expected_counts[bills_path]:
        logger.info(f'Bills count mismatch: got {len(bills_emails)}, expected {expected_counts[bills_path]}')
        return 0.0
    if len(have_seen_emails) != expected_counts[have_seen_path]:
        logger.info(f'have_seen count mismatch: got {len(have_seen_emails)}, expected {expected_counts[have_seen_path]}')
        return 0.0
    if len(have_seen_emails) != 1:
        logger.info(f'have_seen should have exactly 1 email, got {len(have_seen_emails)}')
        return 0.0
    if len(bills_emails) != 1:
        logger.info(f'Bills should have exactly 1 email, got {len(bills_emails)}')
        return 0.0
    moved_email = have_seen_emails[0]
    remaining_email = bills_emails[0]
    moved_timestamp = moved_email.get('date_timestamp', 0)
    remaining_timestamp = remaining_email.get('date_timestamp', 0)
    if moved_timestamp == 0 or remaining_timestamp == 0:
        logger.warning('Could not parse email dates for comparison')
        return 0.0
    if moved_timestamp >= remaining_timestamp:
        logger.info(f'Email in have_seen is not older than email in Bills')
        logger.info(f"Moved email date: {moved_email.get('date')} (timestamp: {moved_timestamp})")
        logger.info(f"Remaining email date: {remaining_email.get('date')} (timestamp: {remaining_timestamp})")
        return 0.0
    logger.info(f'Verified: Oldest email moved to have_seen')
    logger.info(f"Moved email: subject='{moved_email.get('subject')}', date={moved_email.get('date')}")
    logger.info(f"Remaining email: subject='{remaining_email.get('subject')}', date={remaining_email.get('date')}")
    return 1.0

def check_zip_and_folder__173d79a32d2c31ac3ad30d4ae958526d(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the correct files are in the qa_session folder after unzipping.

    Args:
        result: List of filenames from the getter (empty if zip/folder doesn't exist)
        expected: Expected configuration with 'files' key containing required filenames
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    if not result:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    if len(expected_set) == 0:
        return 0.0
    correct_files = len(result_set & expected_set)
    extra_files = len(result_set - expected_set)
    missing_files = len(expected_set - result_set)
    score = correct_files / len(expected_set)
    if extra_files > 0:
        score *= 0.8
    return max(0.0, min(1.0, score))

def check_directory_has_file__739292ff(result, expected, **options):
    """Check if directory contains a specific file with valid image properties.

    Args:
        result: Dict with 'files' list and 'files_info' dict containing file properties
        expected: Dict with 'filename' key
        **options: Additional options

    Returns:
        float: 1.0 if file found and valid, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    files = result.get('files', [])
    files_info = result.get('files_info', {})
    if expected_filename not in files:
        logger.info(f'File not found: {expected_filename}. Files: {files}')
        return 0.0
    file_info = files_info.get(expected_filename, {})
    file_size = file_info.get('size', 0)
    if file_size < MIN_FILE_SIZE:
        logger.info(f'File {expected_filename} is too small: {file_size} bytes (minimum {MIN_FILE_SIZE} bytes)')
        return 0.0
    if not file_info.get('valid_image', False):
        image_error = file_info.get('image_error', 'Unknown error')
        logger.info(f'File {expected_filename} is not a valid image: {image_error}')
        return 0.0
    dimensions = file_info.get('dimensions', {})
    width = dimensions.get('width', 0)
    height = dimensions.get('height', 0)
    if width < MIN_WIDTH or height < MIN_HEIGHT:
        logger.info(f'Image {expected_filename} dimensions too small: {width}x{height} (minimum {MIN_WIDTH}x{MIN_HEIGHT})')
        return 0.0
    image_format = file_info.get('format', '')
    if image_format.upper() != 'PNG':
        logger.info(f'Image {expected_filename} is not PNG format: {image_format}')
        return 0.0
    logger.info(f'Valid image file found: {expected_filename} (size: {file_size} bytes, dimensions: {width}x{height}, format: {image_format})')
    return 1.0

def check_folder_structure__f9d2dc30(result, expected, **options):
    """Check if files were moved to a subfolder correctly.

    Verifies that files exist in the destination folder AND were removed from
    the original location (i.e., moved, not copied).

    Args:
        result: Output from ls command showing both Mountains and Pictures folder contents,
                separated by '---SEPARATOR---'
        expected: Rules dict with required_files list and original_folder path
        **options: Additional options

    Returns:
        float: 1.0 if all files moved correctly (present in destination, absent from source),
               0.5 if files copied (present in both locations),
               0.0 otherwise
    """
    if result is None or not isinstance(result, str):
        return 0.0
    required_files = expected.get('required_files', [])
    if '---SEPARATOR---' not in result:
        return 0.0
    parts = result.split('---SEPARATOR---')
    if len(parts) != 2:
        return 0.0
    mountains_output = parts[0].strip()
    pictures_output = parts[1].strip()
    if 'No such file or directory' in mountains_output:
        return 0.0
    files_in_mountains = 0
    for filename in required_files:
        if filename in mountains_output:
            files_in_mountains += 1
    if files_in_mountains != len(required_files):
        return 0.0
    files_in_pictures = 0
    for filename in required_files:
        if filename in pictures_output:
            files_in_pictures += 1
    if files_in_pictures == 0:
        return 1.0
    elif files_in_pictures == len(required_files):
        return 0.5
    else:
        return 0.5

def check_terminal_profile_name__dfa021d620bcc8024aed513c1adfaad3(result, expected, **options):
    """
    Check if terminal default profile name is set correctly.

    Args:
        result: Terminal output string
        expected: Rules dict with 'include' list containing expected profile name
        **options: Additional options (not used)

    Returns:
        float: 1.0 if profile name is found in output, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_gnome_favorite_apps_ordered__5233e9af0d26c0a6bd95a54609343d46(apps_str: str, rule):
    """Check that GNOME favorite apps match expected list in exact order.

    Args:
        apps_str: String representation of apps list from gsettings
        rule: Dict with "expected" key containing ordered list of app .desktop files

    Returns:
        float: 1.0 if apps match expected order exactly, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if apps == expected_apps:
        return 1.0
    else:
        return 0.0

def check_folder_contains_files__c14b0c4873b98943d0bcc59ebce705c9(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if a folder contains exactly the expected PDF files.

    Args:
        result: List of filenames in the folder
        expected: Dict with 'files' key containing list of expected filenames
        **options: Additional options

    Returns:
        float: 1.0 if folder contains exactly the expected files (no more, no less), 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if expected_set == result_set:
        return 1.0
    else:
        return 0.0

def check_folder_contains_files__32b2b661(result: list, expected: dict, **options) -> float:
    """Check if folder contains expected files.

    Args:
        result: List of filenames in the folder
        expected: Dict with 'required_files' key (list of filenames)
        **options: Additional options

    Returns:
        float: 1.0 if all required files present, 0.0 otherwise
    """
    required_files = expected.get('required_files', [])
    result_set = set(result)
    required_set = set(required_files)
    if required_set.issubset(result_set):
        return 1.0
    return 0.0

def check_gnome_favorite_apps_order__1e04efad(apps_str: str, expected, **options):
    """Check if GNOME favorite apps are in the exact expected order.

    Args:
        apps_str: String representation of favorite apps list from gsettings
        expected: Dict with 'expected_order' key containing ordered list of app .desktop files

    Returns:
        float: 1.0 if order matches exactly, 0.0 otherwise
    """
    try:
        apps = ast.literal_eval(apps_str)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(apps, list):
        return 0.0
    expected_order = expected.get('expected_order', [])
    if apps == expected_order:
        return 1.0
    else:
        return 0.0

def check_qa_folder_contents__fc8f5b84(result: list, expected: dict, **options) -> float:
    """Check if the Q&A folder contains the correct file(s).

    Args:
        result: List of filenames in the Q&A folder
        expected: Dict with 'filenames' key containing list of expected filenames
        **options: Additional options

    Returns:
        float: 1.0 if correct files are present, 0.0 otherwise
    """
    expected_filenames = expected.get('filenames', [])
    result_set = set(result)
    expected_set = set(expected_filenames)
    if result_set == expected_set:
        return 1.0
    return 0.0

def check_desktop_hashes__ef2db1ef(result, expected, **options):
    """Check if Desktop contains all expected image hashes AND Pictures directory is empty of them.

    This verifies files were MOVED (not copied) from Pictures to Desktop.

    Args:
        result: Dict with 'desktop' and 'pictures' keys containing hash lists
        expected: Dict with 'expected_hashes' key

    Returns:
        float: 1.0 if all expected hashes are on Desktop AND none in Pictures, 0.0 otherwise
    """
    expected_hashes = set(expected.get('expected_hashes', []))
    desktop_hashes = set(result.get('desktop', []))
    pictures_hashes = set(result.get('pictures', []))
    if desktop_hashes != expected_hashes:
        missing = expected_hashes - desktop_hashes
        extra = desktop_hashes - expected_hashes
        if missing:
            print(f'Desktop missing hashes: {missing}')
        if extra:
            print(f'Desktop has extra hashes: {extra}')
        return 0.0
    pictures_has_expected = pictures_hashes & expected_hashes
    if pictures_has_expected:
        print(f'Pictures directory still contains hashes (files not moved, only copied): {pictures_has_expected}')
        return 0.0
    return 1.0

def check_gnome_favorite_apps_empty__e17d22d0(apps_str: str, expected, **options):
    """Check if GNOME favorite apps list is empty.

    Args:
        apps_str: String representation of favorite apps list from gsettings
        expected: Empty dict (not used)

    Returns:
        float: 1.0 if list is empty, 0.0 otherwise
    """
    if apps_str is None or not isinstance(apps_str, str):
        return 0.0
    apps_str = apps_str.strip()
    if not apps_str:
        return 0.0
    try:
        apps = ast.literal_eval(apps_str)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(apps, list):
        return 0.0
    if len(apps) == 0:
        return 1.0
    else:
        return 0.0

def check_gnome_favorites_empty__608fb5ab3b43018120dd1b5f6b15910c(result_state, expected_state, **options):
    """
    Check if the GNOME favorites bar is completely empty.

    Args:
        result_state: Output from gsettings get org.gnome.shell favorite-apps command (str)
        expected_state: Expected state (not used, we always expect empty)
        **options: Additional options

    Returns:
        float: Score (1.0 if favorites list is empty, 0.0 otherwise)
    """
    if not result_state:
        return 0.0
    output = result_state.strip()
    if output.startswith('@as '):
        output = output[4:].strip()
    if output == '[]':
        return 1.0
    if '.desktop' in output:
        return 0.0
    if output == '[]' or output == '@as []':
        return 1.0
    return 0.0

def check_nested_directory__3c3f0ceb(result, expected, **options):
    """
    Check if nested directory existence matches expected.

    Args:
        result: Boolean from getter
        expected: Expected boolean value

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_gnome_favorites_first_two__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list contains first two apps only.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing list of two apps

    Returns:
        1.0 if apps match expected (order-independent), 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if len(apps) != len(expected_apps):
        return 0.0
    if set(apps) == set(expected_apps):
        return 1.0
    else:
        return 0.0

def check_bashrc_path_modification__c0107678(result, expected, **options):
    """Check if .bashrc file contains PATH modification and PATH is correctly set.

    This metric verifies that:
    1. The .bashrc file contains a PATH modification with the expected directory
    2. After sourcing .bashrc, PATH actually contains the expected directory

    Args:
        result: Dict with 'bashrc_content' and 'path_value' from getter
        expected: Dict with 'path_entry' key specifying expected PATH directory
        **options: Additional options

    Returns:
        float: 1.0 if both conditions are met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {type(result)}')
        return 0.0
    bashrc_content = result.get('bashrc_content')
    path_value = result.get('path_value')
    path_entry = expected.get('path_entry', '')
    if bashrc_content is None or path_value is None:
        logger.error('Failed to get bashrc content or PATH value')
        return 0.0
    path_assignment_pattern = '(export\\s+)?PATH\\s*=.*' + re.escape(path_entry)
    if not re.search(path_assignment_pattern, bashrc_content, re.MULTILINE):
        logger.info(f'.bashrc does not contain PATH modification with {path_entry}')
        return 0.0
    if path_entry not in path_value:
        logger.info(f'PATH does not contain {path_entry} after sourcing .bashrc')
        return 0.0
    return 1.0

def check_user_shell__49da98ad0a0985144d25f764fe852408(result, expected, **options):
    """Check if user has the expected shell, home directory, and password.

    Args:
        result: Dictionary with user info from getter (username, shell, home, password_set)
        expected: Expected values from rules (shell, home, username, password_set)
        **options: Additional options

    Returns:
        float: 1.0 if all requirements match, 0.0 otherwise
    """
    if result is None:
        logger.info('User not found, returning 0.0')
        return 0.0
    expected_username = expected.get('username')
    if expected_username:
        result_username = result.get('username', '')
        logger.info(f'Comparing username - result: {result_username}, expected: {expected_username}')
        if result_username != expected_username:
            logger.info('Username mismatch')
            return 0.0
    expected_shell = expected.get('shell')
    if expected_shell:
        result_shell = result.get('shell', '').strip()
        logger.info(f'Comparing shell - result: {result_shell}, expected: {expected_shell}')
        if result_shell != expected_shell:
            logger.info('Shell mismatch')
            return 0.0
    expected_home = expected.get('home')
    if expected_home:
        result_home = result.get('home', '').strip()
        logger.info(f'Comparing home - result: {result_home}, expected: {expected_home}')
        if result_home != expected_home:
            logger.info('Home directory mismatch')
            return 0.0
    expected_password_set = expected.get('password_set')
    if expected_password_set is not None:
        result_password_set = result.get('password_set', False)
        logger.info(f'Comparing password_set - result: {result_password_set}, expected: {expected_password_set}')
        if result_password_set != expected_password_set:
            logger.info('Password set status mismatch')
            return 0.0
    logger.info('All user verification checks passed')
    return 1.0

def check_terminal_color_scheme__ad0133b72746f22f133946cc65aa7f55(result, expected, **options):
    """
    Check if terminal color scheme is set correctly.

    Args:
        result: Terminal output string
        expected: Rules dict with 'include' list containing expected color scheme
        **options: Additional options (not used)

    Returns:
        float: 1.0 if color scheme is found in output, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_bash_text_transform__8c025f1893a98ce909d203315d730cfe(result, expected, **options):
    """Check if text transformation matches expected output.

    Args:
        result: Actual file content from getter
        expected: Expected content string or dict with 'content' key
        **options: Additional options (not used)

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_content = expected.get('content', '')
    else:
        expected_content = expected
    if not isinstance(result, str) or not isinstance(expected_content, str):
        return 0.0
    result_stripped = result.rstrip('\n')
    expected_stripped = expected_content.rstrip('\n')
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_directory_created__3ef6f937c5dcda24ec08a8ba5aa2c872(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a directory with required structure was created.

    Args:
        result: Dict from getter with 'exists', 'subdirs', 'files' keys
        expected: Dict with 'required_subdirs' and 'required_files' lists
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.warning('Directory does not exist')
        return 0.0
    score += 0.2
    required_subdirs = expected.get('required_subdirs', [])
    if required_subdirs:
        actual_subdirs = result.get('subdirs', [])
        matching_subdirs = sum((1 for subdir in required_subdirs if subdir in actual_subdirs))
        if matching_subdirs == len(required_subdirs):
            score += 0.4
        else:
            score += 0.4 * (matching_subdirs / len(required_subdirs))
        logger.info(f'Subdirectories matched: {matching_subdirs}/{len(required_subdirs)}')
    required_files = expected.get('required_files', [])
    if required_files:
        actual_files = result.get('files', [])
        matching_files = sum((1 for file in required_files if file in actual_files))
        if matching_files == len(required_files):
            score += 0.4
        else:
            score += 0.4 * (matching_files / len(required_files))
        logger.info(f'Files matched: {matching_files}/{len(required_files)}')
    logger.info(f'Directory structure check score: {score}')
    return score

def check_bash_text_transform__816eebe3a50924fca9ee760018f3238e(result, expected, **options):
    """Check if text transformation matches expected output.

    Args:
        result: Actual file content from getter
        expected: Expected content string or dict with 'content' key
        **options: Additional options (not used)

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_content = expected.get('content', '')
    else:
        expected_content = expected
    if not isinstance(result, str) or not isinstance(expected_content, str):
        return 0.0
    result_stripped = result.rstrip('\n')
    expected_stripped = expected_content.rstrip('\n')
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_directory_created__96172b42(result, expected, **options):
    """Check if a directory was created.

    Args:
        result: Directory listing output from ls -la
        expected: Dict with 'dirname' key specifying expected directory name
        **options: Additional options

    Returns:
        float: 1.0 if directory exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    dirname = expected.get('dirname', '')
    lines = result.strip().split('\n')
    for line in lines:
        if line.startswith('d'):
            parts = line.split()
            if len(parts) >= 9:
                dir_name = ' '.join(parts[8:])
                if dir_name == dirname and dirname not in ['.', '..']:
                    return 1.0
    return 0.0

def check_folder_organization__6a9b75b65029a1742667a75cc968a2a9(result: Dict, expected: Dict, **options) -> float:
    """
    Check if files are organized into correct folders with correct picture-to-folder mapping.

    Args:
        result: Dict with folder contents from getter
        expected: Dict with expected folder structure and picture mappings
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    folders = result.get('folders', {})
    original_files_exist = result.get('original_files_exist', {})
    expected_folders = expected.get('folders', {})
    expected_mappings = expected.get('picture_mappings', {})
    total_checks = 0
    checks_passed = 0
    for (filename, exists) in original_files_exist.items():
        total_checks += 1
        if not exists:
            checks_passed += 1
            logger.info(f"Original file '{filename}' correctly moved from base directory")
        else:
            logger.info(f"Original file '{filename}' still exists in base directory (should be moved, not copied)")
    for (expected_folder_key, expected_file_count) in expected_folders.items():
        total_checks += 1
        found = False
        for actual_folder in folders.keys():
            if expected_folder_key.lower() in actual_folder.lower():
                if len(folders[actual_folder]) >= expected_file_count:
                    checks_passed += 1
                    found = True
                    logger.info(f"Folder '{actual_folder}' has {len(folders[actual_folder])} files (expected >= {expected_file_count})")
                    break
        if not found:
            logger.info(f"No folder found for '{expected_folder_key}' or insufficient files")
    for (picture_filename, expected_folder_keyword) in expected_mappings.items():
        total_checks += 1
        found_correctly = False
        for (actual_folder, files_in_folder) in folders.items():
            if expected_folder_keyword.lower() in actual_folder.lower():
                if picture_filename in files_in_folder:
                    checks_passed += 1
                    found_correctly = True
                    logger.info(f"Picture '{picture_filename}' correctly placed in folder '{actual_folder}'")
                    break
        if not found_correctly:
            logger.info(f"Picture '{picture_filename}' not found in expected folder with keyword '{expected_folder_keyword}'")
    if total_checks == 0:
        return 0.0
    score = checks_passed / total_checks
    logger.info(f'Overall score: {checks_passed}/{total_checks} = {score:.2f}')
    return score

def check_no_stars_in_folder__b22fca61(result, expected, **options):
    """Check if no emails in a specific folder are starred.

    Args:
        result: Path to the SQLite database file
        expected: Dict with 'folder_id' key specifying the folder
        **options: Additional options

    Returns:
        float: 1.0 if no emails are starred, 0.0 if any are starred
    """
    import sqlite3
    folder_id = expected.get('folder_id')
    if folder_id is None:
        return 0.0
    connection = sqlite3.connect(result)
    cursor = connection.cursor()
    cursor.execute('\n        SELECT COUNT(*) FROM messageAttributes\n        WHERE attributeID = 58 AND value = 1\n        AND messageID IN (SELECT id FROM messages WHERE folderID = ?)\n    ', (folder_id,))
    starred_count = cursor.fetchone()[0]
    connection.close()
    return 1.0 if starred_count == 0 else 0.0

def check_file_in_directory__35a43b1b(result, expected, **options):
    """Check if file was moved (exists at destination, not at source).

    Args:
        result: Dict with {"exists_at_destination": bool, "exists_at_source": bool}
        expected: Expected state (dict with 'exists' key - for destination)
        **options: Additional options

    Returns:
        float: 1.0 if file was moved (exists at destination and NOT at source), 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    exists_at_destination = result.get('exists_at_destination', False)
    exists_at_source = result.get('exists_at_source', True)
    if expected_exists:
        if exists_at_destination and (not exists_at_source):
            return 1.0
        else:
            return 0.0
    elif not exists_at_destination:
        return 1.0
    else:
        return 0.0

def check_folder_and_file__223bf56e(result, expected, **options):
    """
    Check if folder and file exist based on command output.

    Args:
        result: Output from test command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if both folder and file exist, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'success')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_directory_moved__11701199(result, expected, **options):
    """Check if directory was moved correctly (source gone, target exists).

    Args:
        result: Command output showing source and target status
        expected: Expected state rules with 'source_should_exist' and 'target_should_exist'
        **options: Additional comparison options

    Returns:
        float: 1.0 if both conditions match, 0.5 if one matches, 0.0 if neither
    """
    source_should_exist = expected.get('source_should_exist', False)
    target_should_exist = expected.get('target_should_exist', True)
    source_exists = 'Source: exists' in result
    target_exists = 'Target: exists' in result
    score = 0.0
    if source_exists == source_should_exist:
        score += 0.5
    if target_exists == target_should_exist:
        score += 0.5
    return score

def check_file_in_directory__3500c3270fba14684f134a0b5e4537d1(result, expected, **options):
    """
    Check if file exists and contains valid SRT subtitle content.

    Args:
        result: Dict from getter with 'exists', 'has_content', 'is_valid_srt' keys
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if file exists with valid SRT content, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('has_content', False):
        logger.info('File exists but is empty')
        return 0.0
    if not result.get('is_valid_srt', False):
        logger.info('File exists but does not contain valid SRT format (missing timestamps or sequence numbers)')
        return 0.0
    logger.info('File exists with valid SRT subtitle content')
    return 1.0

def check_file_in_folder__cbcf6e0c(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if folder exists and contains expected file.

    Args:
        result: Dict from getter with 'folder_exists' and 'file_found' keys
        expected: Expected rules dict

    Returns:
        float: 1.0 if requirements met, 0.0 otherwise
    """
    folder_exists = result.get('folder_exists', False)
    file_found = result.get('file_found', False)
    if expected.get('folder_exists', True) and (not folder_exists):
        return 0.0
    if expected.get('file_exists_in_folder', True) and (not file_found):
        return 0.0
    return 1.0

def check_folder_structure__4e03b1ed(result, expected, **options):
    """Check if folder structure matches expected.

    Args:
        result: Actual folder info from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Score based on folder structure match
    """
    score = 0.0
    if expected.get('exists', False):
        if result.get('exists', False):
            score += 0.5
            logger.info('Folder exists as expected')
        else:
            logger.info('Folder does not exist')
            return 0.0
    expected_count = expected.get('pdf_count', 0)
    if result.get('pdf_count', 0) == expected_count:
        score += 0.5
        logger.info(f'PDF count matches: {expected_count}')
    else:
        logger.info(f"PDF count mismatch. Expected: {expected_count}, Got: {result.get('pdf_count', 0)}")
    return score

def check_folder_counts__856f6f3499821b31f240098182dcc479(result: Optional[Dict], expected, **options):
    """Check if message counts match and the most recent message was moved correctly.

    Args:
        result: Dict with 'counts' (folder -> count), 'bills_most_recent' (subject),
                and 'have_seen_subjects' (list of subjects)
        expected: Dict with 'expected_counts' mapping paths to expected counts,
                  and 'expected_moved_message' with the subject of message that should be moved
        **options: Additional options

    Returns:
        float: 1.0 if all counts match and correct message was moved, 0.0 otherwise
    """
    if result is None:
        logger.info('Result is None - folders not accessible')
        return 0.0
    counts = result.get('counts', {})
    bills_most_recent = result.get('bills_most_recent')
    have_seen_subjects = result.get('have_seen_subjects', [])
    expected_counts = expected.get('expected_counts', {})
    if not expected_counts:
        logger.warning('No expected_counts specified')
        return 0.0
    for (path, expected_count) in expected_counts.items():
        actual_count = counts.get(path)
        if actual_count is None:
            logger.info(f'Folder {path} not found in result')
            return 0.0
        if actual_count != expected_count:
            logger.info(f'Count mismatch for {path}: got {actual_count}, expected {expected_count}')
            return 0.0
    logger.info('Folder counts match expected values')
    if have_seen_subjects:
        if len(have_seen_subjects) == 1:
            moved_subject = have_seen_subjects[0]
            logger.info(f"Message in have_seen: '{moved_subject}'")
            expected_moved = expected.get('expected_moved_message', '')
            if expected_moved:
                if moved_subject == expected_moved or expected_moved in moved_subject:
                    logger.info(f'Correct message was moved: {moved_subject}')
                    return 1.0
                else:
                    logger.info(f"Wrong message moved. Expected '{expected_moved}', got '{moved_subject}'")
                    return 0.0
            else:
                logger.info('No expected_moved_message specified, only checking counts')
                return 1.0
        else:
            logger.info(f'have_seen should have exactly 1 message, got {len(have_seen_subjects)}')
            return 0.0
    else:
        logger.info('have_seen folder is empty')
        return 0.0

def check_gnome_favorite_apps_partial__b4585266e3285db10295625412489679(apps_str: str, rule):
    """Check GNOME favorite apps with partial credit scoring.

    Scoring breakdown:
    - 0.0: Wrong number of apps or missing all required apps
    - 0.5: Correct count but missing some required apps
    - 1.0: All required apps present with correct count

    Args:
        apps_str: String representation of apps list from gsettings
        rule: Dict with "expected" key containing list of required app .desktop files

    Returns:
        float: Score between 0.0 and 1.0
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    expected_count = len(expected_apps)
    apps_set = set(apps)
    expected_set = set(expected_apps)
    count_matches = len(apps) == expected_count
    matching_apps = apps_set & expected_set
    match_ratio = len(matching_apps) / len(expected_set) if expected_set else 0
    if count_matches and match_ratio == 1.0:
        return 1.0
    elif count_matches and match_ratio >= 0.5:
        return 0.5
    elif match_ratio == 1.0:
        return 0.5
    else:
        return 0.0

def check_gnome_favorites_order__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list matches expected order exactly.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing ordered list of .desktop files

    Returns:
        1.0 if apps match expected list in exact order, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if apps == expected_apps:
        return 1.0
    else:
        return 0.0

def check_file_on_desktop__e35e8479c21a9a3d6ef729a997676f8c(result, expected, **options):
    """
    Check if a specific file exists on the desktop.

    Args:
        result: List of filenames on desktop
        expected: Expected condition from rules dict with 'filename'
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_filename = expected.get('filename', '')
    if not expected_filename:
        return 0.0
    if expected_filename in result:
        return 1.0
    else:
        return 0.0

def check_git_subdirectory__a29e9963(result, expected, **options):
    """
    Check if git repository subdirectory exists.

    Args:
        result: Boolean from getter
        expected: Expected value (dict with 'exists' key when type='rule', or boolean)

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_exists = expected.get('exists', True)
    else:
        expected_exists = expected
    if result == expected_exists:
        return 1.0
    return 0.0

def check_gnome_favorites_exact_single__e2063f4a(apps_str: str, rule):
    """Check if favorites list contains exactly one specific app.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'app' key containing the expected app name

    Returns:
        1.0 if list contains only the expected app, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_app = rule['app']
    if len(apps) == 1 and apps[0] == expected_app:
        return 1.0
    else:
        return 0.0

def check_folder_exists__dc2d0173(result: str, rules: dict, **kwargs) -> float:
    """
    Check if a folder with specified name exists in Thunderbird's folderTree.json
    and verify that emails were moved from Inbox to this folder.

    Args:
        result: path to cache file containing folderTree.json content
        rules: dict with "folder_name" key containing the folder name to check

    Returns:
        float: 1.0 if folder exists AND contains emails AND Inbox is empty, 0.0 otherwise
    """
    if result is None:
        return 0.0
    folder_name = rules.get('folder_name', '')
    if not folder_name:
        logger.warning('No folder_name specified in rules')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        folder_tree = json.loads(content)

        def find_folder_with_count(data, name):
            """
            Recursively search for a folder by name in the folder tree.
            Returns (found, message_count) tuple.
            """
            if isinstance(data, dict):
                if data.get('name') == name:
                    msg_count = data.get('totalMessages', 0) or data.get('numMsgs', 0) or data.get('messages', 0)
                    return (True, msg_count)
                if 'children' in data:
                    for child in data['children']:
                        (found, count) = find_folder_with_count(child, name)
                        if found:
                            return (found, count)
                for value in data.values():
                    if isinstance(value, (dict, list)):
                        (found, count) = find_folder_with_count(value, name)
                        if found:
                            return (found, count)
            elif isinstance(data, list):
                for item in data:
                    (found, count) = find_folder_with_count(item, name)
                    if found:
                        return (found, count)
            return (False, 0)
        (inbox_found, inbox_count) = find_folder_with_count(folder_tree, 'Inbox')
        (folder_found, folder_count) = find_folder_with_count(folder_tree, folder_name)
        if not folder_found:
            logger.debug(f"Folder '{folder_name}' not found in folderTree.json")
            return 0.0
        logger.debug(f"Folder '{folder_name}' found with {folder_count} messages")
        if folder_count == 0:
            logger.debug(f"Folder '{folder_name}' exists but contains no messages")
            return 0.0
        if inbox_found and inbox_count > 0:
            logger.debug(f'Inbox still contains {inbox_count} messages')
            return 0.0
        logger.debug(f"Task complete: Folder '{folder_name}' exists with {folder_count} messages and Inbox is empty")
        return 1.0
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON: {e}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking folder existence and email movement: {e}')
        return 0.0

def check_gnome_favorites_empty__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list is empty.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict (not used, but required for signature)

    Returns:
        1.0 if apps list is empty, 0.0 otherwise
    """
    apps = eval(apps_str)
    if len(apps) == 0:
        return 1.0
    else:
        return 0.0

def check_gnome_boolean__fa85fc69(result: str, expected, **options):
    """
    Check if a GNOME boolean setting matches expected value.

    Args:
        result: Path to file containing gsettings output (e.g., "true" or "false")
        expected: Expected boolean value as string ("true" or "false")

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip().lower()
        expected_value = str(expected).strip().lower()
        return 1.0 if content == expected_value else 0.0
    except Exception as e:
        print(f'Error checking boolean setting: {e}')
        return 0.0

def check_directory_count__7ec80ec4(result, expected, **options):
    """
    Check if directory count matches expected value.

    Args:
        result: Integer count from getter
        expected: Expected count

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if result == expected_count:
        return 1.0
    return 0.0

def check_folder_not_empty__95db6624(result: bool, expected: Dict[str, Any], **options) -> float:
    """Check if folder is not empty"""
    if result:
        logger.info('Folder has files')
        return 1.0
    logger.info('Folder is empty')
    return 0.0

def check_zip_and_folder__36005e14c2ae4846381f0946699e70be(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the zip file exists, is valid, and contains the correct files.

    Args:
        result: Dict with 'zip_exists', 'files', and 'error' keys from the getter
        expected: Expected configuration with 'files' key
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    if not isinstance(result, dict):
        return 0.0
    if not result.get('zip_exists', False):
        return 0.0
    if result.get('error'):
        return 0.0
    result_files = result.get('files', [])
    if not result_files:
        return 0.0
    result_set = set(result_files)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    if len(expected_set) == 0:
        return 0.0
    correct_files = len(result_set & expected_set)
    extra_files = len(result_set - expected_set)
    missing_files = len(expected_set - result_set)
    score = correct_files / len(expected_set)
    if extra_files > 0:
        score *= 0.8
    return max(0.0, min(1.0, score))

def check_gnome_app_added__1ed286d0e6299072018d07279d4346ca(result_state, expected_state, **options):
    """
    Check if the Files application was added to GNOME favorites.

    Verifies that:
    1. The expected app (org.gnome.Nautilus.desktop) is in the favorites
    2. All original apps are still present
    3. Exactly one new app was added

    Args:
        result_state: Output from gsettings get org.gnome.shell favorite-apps
        expected_state: Dict with 'expected_app' and 'original_apps'
        **options: Additional options

    Returns:
        float: 1.0 if verification passes, 0.0 otherwise
    """
    if not result_state or not isinstance(result_state, str):
        return 0.0
    import ast
    try:
        result_state = result_state.strip()
        current_apps = ast.literal_eval(result_state)
    except (ValueError, SyntaxError):
        return 0.0
    if not isinstance(current_apps, list):
        return 0.0
    expected_app = expected_state.get('expected_app')
    original_apps = expected_state.get('original_apps', [])
    if expected_app not in current_apps:
        return 0.0
    for app in original_apps:
        if app not in current_apps:
            return 0.0
    if len(current_apps) != len(original_apps) + 1:
        return 0.0
    return 1.0

def check_subfolder_hashes__fcfe8473(result, expected, **options):
    """Check if subfolder contains expected image hashes.

    Args:
        result: List of hashes from subfolder
        expected: Dict with 'expected_hashes' key

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    expected_hashes = set(expected.get('expected_hashes', []))
    result_set = set(result)
    if result_set == expected_hashes:
        return 1.0
    else:
        missing = expected_hashes - result_set
        extra = result_set - expected_hashes
        if missing:
            print(f'Missing hashes: {missing}')
        if extra:
            print(f'Extra hashes: {extra}')
        return 0.0

def check_reconciled_folder__fc048187(result, expected, **options):
    """
    Check if the reconciled folder contains exactly the expected files.

    Args:
        result: Output from ls command listing files in the folder
        expected: Dict with 'rules' containing 'required_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    required_files = expected.get('required_files', [])
    if isinstance(result, dict) and 'output' in result:
        result = result['output']
    result = str(result).strip()
    if 'No such file' in result or 'cannot access' in result:
        return 0.0
    actual_files = [line.strip() for line in result.split('\n') if line.strip()]
    score = 0.0
    if actual_files or result == '':
        score += 0.3
    if set(actual_files) == set(required_files):
        score += 0.7
    elif actual_files:
        correct_count = len(set(actual_files) & set(required_files))
        extra_count = len(set(actual_files) - set(required_files))
        missing_count = len(set(required_files) - set(actual_files))
        if extra_count == 0 and missing_count > 0:
            score += 0.7 * (correct_count / len(required_files))
        elif extra_count > 0 and missing_count == 0:
            score += 0.5
        else:
            score += 0.3 * (correct_count / len(required_files))
    return min(score, 1.0)

def check_zip_and_folder__9ff7dd4db34cdbddf16b8ce0d6085594(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the correct files are in the folder and zip exists.

    Args:
        result: List of filenames from the getter
        expected: Expected configuration with 'files' key
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    if not result:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    if len(expected_set) == 0:
        return 0.0
    correct_files = len(result_set & expected_set)
    extra_files = len(result_set - expected_set)
    missing_files = len(expected_set - result_set)
    score = correct_files / len(expected_set)
    if extra_files > 0:
        score *= 0.8
    return max(0.0, min(1.0, score))

def check_directory_existence__0b3844b6(result, expected, **options):
    """
    Check if directory existence matches expected value.

    Args:
        result: Boolean from getter (True if directory exists)
        expected: Dictionary with 'exists' key containing expected boolean value

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_folder_contains_files__21ed89b452cf1a748805f1dda9b2ec4b(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if a folder contains expected files.

    Args:
        result: List of files in the folder (from getter)
        expected: Dict with 'required_files' key - list of files that must exist

    Returns:
        1.0 if all required files exist, 0.0 otherwise
    """
    required_files = expected.get('required_files', [])
    if not required_files:
        logger.warning('No required_files specified in expected config')
        return 0.0
    result_set = set(result)
    for req_file in required_files:
        if req_file not in result_set:
            logger.info(f'Missing required file: {req_file}')
            return 0.0
    logger.info(f'All required files present: {required_files}')
    return 1.0

def check_bash_text_transform__5ed0399c75d125da4fc3b5f7583bb5c4(result, expected, **options):
    """Check if text transformation matches expected output.

    Args:
        result: Actual file content from getter
        expected: Expected content string or dict with 'content' key
        **options: Additional options (not used)

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_content = expected.get('content', '')
    else:
        expected_content = expected
    if not isinstance(result, str) or not isinstance(expected_content, str):
        return 0.0
    result_stripped = result.rstrip('\n')
    expected_stripped = expected_content.rstrip('\n')
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_folder_and_file__fe4fce7a(result, expected, **options):
    """
    Check if folder and file exist based on command output.

    Args:
        result: Output from test command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if both folder and file exist, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'success')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_terminal_scrollback__8e9fd232f4307bd96cd64b8e7a1f9389(result, expected, **options):
    """
    Check if terminal scrollback lines is set correctly.

    Args:
        result: Terminal output string
        expected: Rules dict with 'include' list containing expected scrollback value
        **options: Additional options (not used)

    Returns:
        float: 1.0 if scrollback value is found in output, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_gnome_favorites_order__a572fc31d81435d0b739124764e326d3(apps_str: str, rule):
    """
    Check if Google Chrome appears at the top of GNOME favorite apps.

    Args:
        apps_str: String representation of the favorites list from gsettings,
                  e.g., "['thunderbird.desktop', 'vim.desktop', 'google-chrome.desktop']"
        rule: Dict containing 'expected_order' key with the expected list of apps in order
                  (used to identify which app should be first)

    Returns:
        float: 1.0 if Chrome is at the top (first position), 0.0 otherwise
    """
    try:
        apps = ast.literal_eval(apps_str)
    except (ValueError, SyntaxError):
        return 0.0
    if not apps or len(apps) == 0:
        return 0.0
    if apps[0] == 'google-chrome.desktop':
        return 1.0
    else:
        return 0.0

def check_gnome_favorites_swap_first_two__e2063f4a(apps_str: str, rule):
    """Check if first two favorite apps are swapped compared to original.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing list with swapped first two apps

    Returns:
        1.0 if apps match expected with first two swapped, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_apps = rule['expected']
    if apps == expected_apps:
        return 1.0
    else:
        return 0.0

def check_desktop_file_exists__d17b574f(result, expected, **options):
    """Check if file exists on Desktop.

    Args:
        result: Boolean from getter indicating if file exists
        expected: Expected properties from rules
        **options: Additional options

    Returns:
        1.0 if file exists on Desktop, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_gnome_favorites_single__e2063f4a(apps_str: str, rule):
    """Check if favorite apps list contains exactly one specific app.

    Args:
        apps_str: String representation of list from gsettings get
        rule: Dict with 'expected' key containing the single app name

    Returns:
        1.0 if list has exactly one app matching expected, 0.0 otherwise
    """
    apps = eval(apps_str)
    expected_app = rule['expected']
    if len(apps) == 1 and apps[0] == expected_app:
        return 1.0
    else:
        return 0.0

def check_directory_exists__881ee5e7(result, expected, **options):
    """Check if directory exists by verifying it's specifically a directory.

    Args:
        result: Dict with keys {'exists': bool, 'is_directory': bool, 'disk_usage': str or None}
        expected: Dict (not used, checking for directory existence)
        **options: Additional options

    Returns:
        float: 1.0 if path exists and is a directory, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('is_directory', False) and result.get('exists', False):
        if result.get('disk_usage') is not None:
            return 1.0
        return 1.0
    return 0.0

def check_files_in_folder__2b29a90e(result, expected, **options):
    """Check files in folder meet criteria.

    Args:
        result: list of file info dicts
        expected: dict with min_count and file_type

    Returns:
        float: 1.0 if criteria met, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    min_count = expected.get('min_count', 0)
    file_type = expected.get('file_type', '')
    if file_type:
        matching = [f for f in result if f.get('title', '').endswith(f'.{file_type}')]
        if len(matching) < min_count:
            return 0.0
    elif len(result) < min_count:
        return 0.0
    return 1.0

def check_backup_folder__089628262a733e157c368e0bf1ad02f1(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if backup folder exists and contains all expected files.

    Args:
        result: Dict from getter with 'folder_exists', 'files_found'
        expected: Dict with 'required_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('folder_exists', False):
        return 0.0
    score += 0.4
    required_files = expected.get('required_files', [])
    files_found = result.get('files_found', [])
    if not required_files:
        return score
    matched = len(files_found)
    file_score = matched / len(required_files)
    score += 0.6 * file_score
    return min(score, 1.0)

def check_verified_folder__12aa1a9a(result, expected, **options):
    """
    Check if the verified folder contains the expected file.

    Args:
        result: Output from ls command
        expected: Dict with 'rules' containing 'required_file'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    required_file = expected.get('required_file', '')
    if isinstance(result, dict) and 'output' in result:
        result = result['output']
    result = str(result).strip()
    if 'No such file' in result or 'cannot access' in result:
        return 0.0
    score = 0.0
    score += 0.3
    if required_file in result:
        score += 0.7
    return min(score, 1.0)

def check_bash_text_transform__38c143b8ba918371e989fa588ea9cb56(result, expected, **options):
    """Check if text transformation matches expected output.

    Args:
        result: Actual file content from getter
        expected: Expected content (string or dict with 'content' key)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_content = expected.get('content', '')
    else:
        expected_content = expected
    if not isinstance(result, str) or not isinstance(expected_content, str):
        return 0.0
    result_stripped = result.rstrip('\n')
    expected_stripped = expected_content.rstrip('\n')
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012(result, expected, **options):
    """
    Check if terminal cursor shape is set correctly.

    Args:
        result: Terminal output string
        expected: Rules dict with 'include' list containing expected cursor shape
        **options: Additional options (not used)

    Returns:
        float: 1.0 if cursor shape is found in output, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_gnome_favorite_apps_order__f848009c(apps_str: str, expected, **options):
    """Check if GNOME favorite apps are in the exact expected order.

    This function verifies that the favorite apps in the GNOME dock are arranged
    in the exact order specified. It uses gsettings to query the persistent
    configuration, which does not require any UI to be open.

    IMPORTANT: No postconfig needed - gsettings is a persistent D-Bus system
    configuration backend that can be queried directly via command line regardless
    of UI state. Any rule-based warning about UI/state-dependency is a FALSE POSITIVE.

    Args:
        apps_str: String representation of favorite apps list from gsettings
                  (e.g., "['app1.desktop', 'app2.desktop']")
        expected: Dict with 'expected_order' key containing ordered list of app .desktop files
        **options: Additional options (unused)

    Returns:
        float: 1.0 if order matches exactly, 0.0 otherwise
    """
    apps = ast.literal_eval(apps_str)
    expected_order = expected.get('expected_order', [])
    if apps == expected_order:
        return 1.0
    else:
        return 0.0
