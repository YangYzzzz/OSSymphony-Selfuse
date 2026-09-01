"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_vscode_tab_size__6ee0182d', 'check_vscode_theme__f8298be8', 'check_vscode_open_file__ddd2693d', 'check_json_keybindings__afd09745', 'check_workspace_folders__5eef46805c2f01e88d41b2c3c5e6f09a', 'check_vscode_terminal_open__64ba4345', 'check_vscode_contains_text__a5163ad1', 'check_json_keybindings__53bf24f6', 'check_vscode_tabs_not_contain__32a2204c', 'check_vscode_explorer__71b2b2791188a1049bb680b7c56ffb2d', 'check_json_keybindings__bf7238c9', 'check_workspace_folders__5ec340b8da9a9ba75d41c5d950fd9065', 'check_json_keybindings__4e0013e9', 'check_vscode_workspace_folders__6ed0a554', 'check_vscode_search__0da38d8e', 'check_workspace_folder_order__e277e2d9', 'check_vscode_indentation__b5e3afd36096fa6a3d39a615d4795ccf', 'check_vscode_launch_json__6ee0182d', 'check_vscode_setting__052d7bff3ecc5844bf94710602cff931', 'check_vscode_indentation__59124b969e280d0086072fa0af69d61b', 'check_json_keybindings__97411736', 'check_vscode_workspace_file__56d366f94894cbfe50f3e560644c1cc7', 'check_json_keybindings__71ff2c93', 'check_vscode_settings_key__a74ebb334721e27ebd0b52511b5b5b07', 'check_keybinding_new_terminal__268ff7f264fd39b73aa8c06f129daf92', 'check_vscode_autosave__c5efe7eb60d4b82c4d217326823a4934', 'check_vscode_nodemod_exclude__bdfa2ae4136b63f287ad1f185959812f', 'check_vscode_setting__d0166485dd9dee18fc251c7dbffffc79', 'check_vscode_not_contains__02be85da', 'check_vscode_wordwrap__ca4beaed446890c3b888d4a40b939831', 'check_vscode_variable_rename__09b11b45', 'check_workspace_folder_count__8a8ecf82', 'check_vscode_settings_values__51833a76bce91777752996df0b22c33d', 'check_vscode_line_content__da119c1e', 'check_json_keybindings__c92c626f', 'check_vscode_formatonsave__4df3f2e6', 'check_vscode_function_name__75a1331b', 'check_vscode_indentation__06f73a877cbecd1e2a247622c3136810', 'check_vscode_setting__4269362e9e26c329da428a21ddabb7ea', 'check_vscode_contains_text__2fa9f229', 'check_vscode_open_file__005f6ac9', 'check_json_keybindings__c38e7539', 'check_json_keybindings__3742c48e', 'check_vscode_wordwrap__ff65f872', 'check_vscode_setting__37a81987aef2f239d369b6635778b9ec', 'check_vscode_tabsize__a6b4939a2d71540bebc75ede5d00109a', 'check_vscode_config_files__17da0621', 'check_workspace_folder_order__40d497fa35fd103d8064472e794b77f6', 'check_json_keybindings__8bf59803', 'check_workspace_folders__328ca95134872b7162fe90bdb76e5fde', 'check_vscode_open_file__149df8b642d00f3a86d140b4582f20cc', 'check_json_keybindings__9f7d30da', 'check_keybinding_editor_to_terminal__d498cab59d97f4e17938744aa091cf0b', 'check_multiple_workspaces__f54bf781c32dea93b80d6edc50522687', 'check_json_keybindings__ec7a7a17', 'check_json_keybindings__45d3929c', 'check_workspace_has_folder__33d8878a', 'check_vscode_indentation__219e3148d77134036e282c6dd4b41d12', 'check_json_keybindings__a807c39e', 'check_workspace_folders__ff3abf45ba7fc3112a6dd2bb6da6358b', 'check_json_keybindings__c2ab0a78', 'check_json_keybindings__093ec527', 'check_json_keybindings__940d9534', 'check_vscode_setting__a46d28041c0fe91cc8b08a2032679781', 'check_vscode_markdown_files__77849053', 'check_vscode_workspace__6efce4eb1d6563b7b443059d330de6aa', 'check_workspace_folders__9f48007a04e5b61abc6f8fde0b255a1d', 'check_json_keybindings__76d7a1f4', 'check_vscode_autosave__aec709e9', 'check_vscode_format_on_paste__6ee0182d', 'check_vscode_open_file__c60120ed87e101802d8b786344162dfa', 'check_vscode_workspace_files__4ab817c3', 'check_vscode_setting__b07969a1', 'check_vscode_line_content__9d7d44f7', 'check_json_keybindings__e8fad444', 'check_vscode_git_exclude__3826e4daac10a542e57e7cc4a2b6b258', 'check_vscode_file_count__7c2e5da9', 'check_vscode_tabsize__44a94327', 'check_vscode_indentation__57096d124fd446ca6cba7d1316aac1bd', 'check_vscode_open_file__e925f44c', 'check_keybinding_sidebar_toggle__9bbee30a9c2c82b51b389394e0233d5e', 'check_json_keybindings__9c5a8886', 'check_vscode_settings_autosave__06d1442bba585f1fb8ea50a4956df674', 'check_workspace_folders', 'check_keybinding_terminal_toggle__326285452dea673d67652e18309c6e13', 'check_workspace_excludes_folder__2068a9c3', 'check_vscode_contains_text__39395653', 'check_vscode_setting__232783f17230c635388a2d96b7097023', 'check_keybinding_reopen_editor__252b583b8604f0bd370b6e7ab984580e', 'check_vscode_format_on_save__7a482b52']

def check_vscode_tab_size__6ee0182d(actual: str, rules: dict, **options) -> float:
    """
    Check if VSCode settings.json has the correct tab size setting.

    Args:
        actual (str): path to settings.json file
        rules (dict): expected configuration rules with tab_size

    Returns:
        float: score between 0.0 and 1.0
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Failed to load settings.json: {e}')
        return 0.0
    expected_tab_size = rules.get('tab_size', 2)
    score = 0.0
    if 'editor.tabSize' in data:
        if data['editor.tabSize'] == expected_tab_size:
            score = 1.0
    return score

def check_vscode_theme__f8298be8(actual: str, rules: Dict, **options) -> float:
    """
    Check if VS Code theme is set to the expected value.

    Args:
        actual (str): Path to the settings.json file
        rules (Dict): Expected rules with 'expected_theme' key

    Returns:
        float: 1.0 if theme matches, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expected_theme = rules.get('expected_theme', '')
    actual_theme = data.get('workbench.colorTheme', '')
    if actual_theme == expected_theme:
        return 1.0
    return 0.0

def check_vscode_open_file__ddd2693d(actual: str, rules: dict, **options) -> float:
    """Check if the expected file is open in VSCode.

    Args:
        actual: Path to the file containing the open file info
        rules: Dict with 'expected' key containing the expected file name/path
        **options: Additional options

    Returns:
        1.0 if the expected file is open, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read().strip()
        expected = rules['expected']
        if expected in content:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking open file: {e}')
        return 0.0

def check_json_keybindings__afd09745(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (from vm_file getter)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    file_path = actual
    data = None
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
    except:
        try:
            with open(file_path, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            return 0.0
    if data is None or not isinstance(data, list):
        return 0.0
    expected_binding = expected.get('expected')
    for binding in data:
        if binding.get('key') == expected_binding.get('key') and binding.get('command') == expected_binding.get('command'):
            return 1.0
    return 0.0

def check_workspace_folders__5eef46805c2f01e88d41b2c3c5e6f09a(actual: str, expected: Dict, **options) -> float:
    """
    Check if workspace file contains expected folders.

    Args:
        actual (str): path to workspace file
        expected (Dict): expected configuration with 'folders' key

    Returns:
        float: 1.0 if folders match, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expected_folders = expected.get('folders', [])
    actual_folders = data.get('folders', [])
    if len(expected_folders) != len(actual_folders):
        return 0.0
    expected_paths = {folder['path'] for folder in expected_folders}
    actual_paths = {folder['path'] for folder in actual_folders}
    if expected_paths == actual_paths:
        return 1.0
    return 0.0

def check_vscode_terminal_open__64ba4345(actual: str, expected: Dict, **options) -> float:
    """Check if the terminal is open in VS Code.

    Args:
        actual (str): path to result text file containing terminal status
        expected (Dict): expected dict with 'is_open' key
        **options: Additional options

    Return:
        float: the score (1.0 if terminal status matches, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
    except Exception:
        return 0.0
    is_open = expected.get('is_open', True)
    terminal_open = 'open' in actual_text.lower() or 'true' in actual_text.lower()
    if terminal_open == is_open:
        return 1.0
    return 0.0

def check_vscode_contains_text__a5163ad1(result, expected, **options):
    """Check if file contains expected text strings and does not contain old variable names."""
    import re
    import ast
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        expected_strings = expected.get('contains', [])
        for s in expected_strings:
            if s not in content:
                return 0.0
        old_var_names = ['n', 'j', 'i']
        try:
            tree = ast.parse(content)
            var_names = set()
            for node in ast.walk(tree):
                if isinstance(node, ast.Name):
                    var_names.add(node.id)
            for old_var in old_var_names:
                if old_var in var_names:
                    return 0.0
        except SyntaxError:
            for old_var in old_var_names:
                pattern = '\\b' + re.escape(old_var) + '\\b'
                if re.search(pattern, content):
                    return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_json_keybindings__53bf24f6(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 5: Disable Alt+Up shortcut for editor.action.moveLinesUpAction

    Args:
        actual (str): path to keybindings.json file (from vm_file result)
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if not isinstance(data, list):
            raise ValueError('Expected JSON array')
    except:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
            if not isinstance(data, list):
                return 0.0
        except:
            return 0.0
    expected_entry = expected['expected']
    if expected_entry in data:
        return 1.0
    else:
        return 0.0

def check_vscode_tabs_not_contain__32a2204c(actual: str, expected: Dict, **options) -> float:
    """Check that a specific file is NOT in the open tabs.

    Args:
        actual (str): path to result text file containing open tabs info
        expected (Dict): expected dict with 'excluded_file' key
        **options: Additional options

    Return:
        float: the score (1.0 if file is NOT open, 0.0 if it is open)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
    except Exception:
        return 0.0
    excluded_file = expected.get('excluded_file', '')
    if excluded_file not in actual_text:
        return 1.0
    return 0.0

def check_vscode_explorer__71b2b2791188a1049bb680b7c56ffb2d(actual: str, rules: Dict, **options) -> float:
    """Check if the expected files are visible in the VSCode explorer.

    Args:
        actual: Path to the file containing the explorer content
        rules: Dict with 'expected_files' list of filenames to check
        **options: Additional options

    Returns:
        float: Score based on how many expected files are found
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
        expected_files = rules.get('expected_files', [])
        if not expected_files:
            return 0.0
        found_count = 0
        for filename in expected_files:
            if filename in actual_text:
                found_count += 1
        return found_count / len(expected_files)
    except Exception as e:
        logger.error(f'Error checking explorer: {e}')
        return 0.0

def check_json_keybindings__bf7238c9(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    file_path = actual
    if not os.path.exists(file_path):
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    data = None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    if data is None:
        return 0.0
    expected_binding = expected.get('expected')
    for binding in data:
        if binding.get('key') == expected_binding.get('key') and binding.get('command') == expected_binding.get('command') and (binding.get('when') == expected_binding.get('when')):
            return 1.0
    return 0.0

def check_workspace_folders__5ec340b8da9a9ba75d41c5d950fd9065(actual: str, expected: Dict, **options) -> float:
    """
    Check if workspace file contains expected folders.

    Args:
        actual (str): path to workspace file
        expected (Dict): expected configuration with 'folders' key

    Returns:
        float: 1.0 if folders match, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expected_folders = expected.get('folders', [])
    actual_folders = data.get('folders', [])
    if len(expected_folders) != len(actual_folders):
        return 0.0
    expected_paths = {folder['path'] for folder in expected_folders}
    actual_paths = {folder['path'] for folder in actual_folders}
    if expected_paths == actual_paths:
        return 1.0
    return 0.0

def check_json_keybindings__4e0013e9(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (from vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    file_path = actual
    try:
        with open(file_path, 'r') as f:
            data = json.load(f)
        if type(data) == list:
            keybindings = data
        else:
            return 0.0
    except:
        try:
            with open(file_path, 'r') as f:
                f.readline()
                data = json.load(f)
            if type(data) == list:
                keybindings = data
            else:
                return 0.0
        except:
            return 0.0
    expected_binding = expected.get('expected')
    for binding in keybindings:
        if binding.get('key') == expected_binding.get('key') and binding.get('command') == expected_binding.get('command'):
            return 1.0
    return 0.0

def check_vscode_workspace_folders__6ed0a554(result_state: Optional[Dict], expected_state: Dict, **options) -> float:
    """
    Check if VSCode workspace contains the required folders.

    This metric verifies that specific folders are present in the workspace,
    with support for both absolute and relative path formats.

    Args:
        result_state: Dictionary containing the workspace JSON data from getter
        expected_state: Dict with 'required_folders' key containing list of folder paths to check
        **options: Additional options

    Returns:
        float: Score 1.0 if all required folders are present, 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Result state is None')
        return 0.0
    if not isinstance(result_state, dict):
        logger.error(f'Result state is not a dict: {type(result_state)}')
        return 0.0
    workspace_folders = result_state.get('folders', [])
    if not isinstance(workspace_folders, list):
        logger.error(f'Workspace folders is not a list: {type(workspace_folders)}')
        return 0.0
    required_folders = expected_state.get('required_folders', [])
    if not isinstance(required_folders, list):
        logger.error(f'Required folders is not a list: {type(required_folders)}')
        return 0.0

    def normalize_path(path_str: str) -> str:
        """
        Normalize path to handle both absolute and relative formats.
        Converts to absolute format for comparison.
        """
        path_str = path_str.strip()
        if path_str.startswith('/'):
            return path_str
        return os.path.join('/home/user', path_str)
    workspace_paths = set()
    for folder in workspace_folders:
        if isinstance(folder, dict) and 'path' in folder:
            normalized = normalize_path(folder['path'])
            workspace_paths.add(normalized)
            logger.debug(f"Found workspace folder: {folder['path']} -> {normalized}")
    required_paths = set()
    for folder_path in required_folders:
        normalized = normalize_path(folder_path)
        required_paths.add(normalized)
        logger.debug(f'Required folder: {folder_path} -> {normalized}')
    missing_folders = required_paths - workspace_paths
    if missing_folders:
        logger.info(f'Missing required folders: {missing_folders}')
        logger.info(f'Workspace has: {workspace_paths}')
        return 0.0
    logger.info(f'All required folders present: {required_paths}')
    return 1.0

def check_vscode_search__0da38d8e(actual: str, expected: Dict, **options) -> float:
    """Check if the search was performed with the correct query and returned results.

    Args:
        actual (str): path to result text file containing search results
        expected (Dict): expected dict with 'search_query' key
        **options: Additional options

    Return:
        float: the score (1.0 if search was executed with results, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
    except Exception:
        return 0.0
    search_query = expected.get('search_query', '')
    if not actual_text:
        return 0.0
    if search_query.lower() not in actual_text.lower():
        return 0.0
    has_results_indicators = any(['main.py' in actual_text, '.py' in actual_text, '/home/user/project' in actual_text, 'match' in actual_text.lower(), 'result' in actual_text.lower(), ':' in actual_text])
    if has_results_indicators:
        return 1.0
    return 0.0

def check_workspace_folder_order__e277e2d9(result: str, expected, **options) -> float:
    """
    Check if workspace folders are in the expected order.

    Args:
        result (str): path to workspace file
        expected: dict containing "folder_paths" key with list of folder paths in expected order
        **options: additional options

    Returns:
        float: 1.0 if folders are in correct order, 0.0 otherwise
    """
    if not result:
        logger.warning('Result file path is None, returning 0.0')
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Error loading JSON from {result}: {e}')
        return 0.0
    folders = data.get('folders', [])
    actual_paths = [folder.get('path', '') for folder in folders]
    expected_paths = expected.get('folder_paths', [])
    actual_basenames = [os.path.basename(path.rstrip('/')) for path in actual_paths]
    expected_basenames = [os.path.basename(path.rstrip('/')) for path in expected_paths]
    if actual_basenames == expected_basenames:
        return 1.0
    else:
        logger.debug(f'Folder order mismatch: expected {expected_basenames}, got {actual_basenames}')
        return 0.0

def check_vscode_indentation__b5e3afd36096fa6a3d39a615d4795ccf(result, expected, **options):
    """Compare file content with expected text.

    Args:
        result: Actual file content (string)
        expected: Expected content (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if content matches exactly, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_content = expected.get('content', '')
    if not expected_content:
        return 0.0
    if result == expected_content:
        return 1.0
    return 0.0

def check_vscode_launch_json__6ee0182d(actual: str, rules: dict, **options) -> float:
    """
    Check if launch.json exists and has correct Python debug configuration.

    Args:
        actual (str): path to launch.json file
        rules (dict): expected configuration rules

    Returns:
        float: score between 0.0 and 1.0
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Failed to load launch.json: {e}')
        return 0.0
    score = 0.0
    if 'configurations' not in data or not isinstance(data['configurations'], list):
        return 0.0
    if len(data['configurations']) == 0:
        return 0.0
    score += 0.5
    has_python_config = False
    for config in data['configurations']:
        if config.get('type') == 'python':
            has_python_config = True
            break
    if has_python_config:
        score += 0.5
    return score

def check_vscode_setting__052d7bff3ecc5844bf94710602cff931(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pair.

    Args:
        result: Dictionary from getter containing actual VS Code settings
        expected: Dictionary with 'key' and 'value' to check (from rules)
        **options: Additional options

    Returns:
        1.0 if the setting key has the expected value, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    try:
        if float(actual_value) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_vscode_indentation__59124b969e280d0086072fa0af69d61b(result, expected, **options):
    """Compare file content with expected text.

    Args:
        result: Actual file content (string)
        expected: Expected content (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if content matches exactly, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_content = expected.get('content', '')
    if not expected_content:
        return 0.0
    if result == expected_content:
        return 1.0
    return 0.0

def check_json_keybindings__97411736(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 0: Remove Ctrl+D shortcut for editor.action.addSelectionToNextFindMatch

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            return 0.0
    if data is None or not isinstance(data, list):
        return 0.0
    expected_entry = expected['expected']
    if expected_entry in data:
        return 1.0
    else:
        return 0.0

def check_vscode_workspace_file__56d366f94894cbfe50f3e560644c1cc7(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the workspace is actually opened in VSCode.

    Args:
        result: Dict from getter with 'vscode_running', 'workspace_opened', 'file_exists', 'valid_json', 'folders' keys
        expected: Dict with 'workspace_path' and 'expected_folders' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('vscode_running', False):
        score += 0.2
    if result.get('workspace_opened', False):
        score += 0.4
    if result.get('file_exists', False):
        score += 0.2
    if result.get('valid_json', False):
        score += 0.1
        expected_folders = expected.get('expected_folders', [])
        result_folders = result.get('folders', [])
        if expected_folders:
            expected_paths = {f.get('path', '') if isinstance(f, dict) else str(f) for f in expected_folders}
            result_paths = {f.get('path', '') if isinstance(f, dict) else str(f) for f in result_folders}
            if expected_paths.issubset(result_paths):
                score += 0.1
    return score

def check_json_keybindings__71ff2c93(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 8: Remove Ctrl+Shift+[ shortcut for editor.fold

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if not actual or not os.path.exists(actual):
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if data is not None and isinstance(data, list):
            expected_entry = expected['expected']
            return 1.0 if expected_entry in data else 0.0
    except (json.JSONDecodeError, Exception):
        pass
    try:
        with open(actual, 'r') as f:
            f.readline()
            data = json.load(f)
        if data is not None and isinstance(data, list):
            expected_entry = expected['expected']
            return 1.0 if expected_entry in data else 0.0
    except (json.JSONDecodeError, Exception):
        pass
    return 0.0

def check_vscode_settings_key__a74ebb334721e27ebd0b52511b5b5b07(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if a specific key exists in VS Code settings with the expected value.

    Args:
        result: Dict containing the settings from getter
        expected: Dict with 'key' and 'value' fields specifying what to check
        **options: Additional options (not used)

    Returns:
        float: 1.0 if key exists with expected value, 0.0 otherwise
    """
    if not result:
        return 0.0
    key = expected.get('key')
    value = expected.get('value')
    if key is None:
        logger.error('Expected key not provided')
        return 0.0
    if key not in result:
        return 0.0
    if result[key] == value:
        return 1.0
    return 0.0

def check_keybinding_new_terminal__268ff7f264fd39b73aa8c06f129daf92(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding is configured in VSCode keybindings.json.
    This variation checks for new terminal creation keybinding.

    Args:
        actual (str): path to result keybindings.json file
        expected (Dict): expected dict with 'expected' key containing the keybinding config

    Return:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    else:
        return 0.0
    expected_binding = expected['expected']
    if expected_binding in data:
        return 1.0
    else:
        return 0.0

def check_vscode_autosave__c5efe7eb60d4b82c4d217326823a4934(actual: str, expected: Dict, **options) -> float:
    """Check if VS Code settings.json has auto-save configured correctly.

    Args:
        actual: Path to settings.json file
        expected: Expected configuration rules (from expected.rules)
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel or actual_value != value:
            return 0.0
    return 1.0

def check_vscode_nodemod_exclude__bdfa2ae4136b63f287ad1f185959812f(actual: str, expected: Dict, **options) -> float:
    """Check if VS Code settings.json has node_modules folder exclusion configured.

    Args:
        actual: Path to settings.json file
        expected: Expected configuration rules (from expected.rules)
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel or actual_value != value:
            return 0.0
    return 1.0

def check_vscode_setting__d0166485dd9dee18fc251c7dbffffc79(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pair.

    Args:
        result: Dictionary from getter containing actual VS Code settings
        expected: Dictionary with 'key' and 'value' to check (from rules)
        **options: Additional options

    Returns:
        1.0 if the setting key has the expected value, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    try:
        if float(actual_value) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_vscode_not_contains__02be85da(result, expected, **options):
    """Check if specific lines containing the pattern were deleted.

    This checks for the multi-line pattern constructed from expected['not_contains']
    that should be removed, rather than checking for substring absence globally
    (which would incorrectly fail if the strings appear elsewhere in the file).

    Args:
        result: Path to the file to check
        expected: Dict containing 'not_contains' list of strings that should not appear together
        **options: Additional options

    Returns:
        float: 1.0 if pattern is not found, 0.0 otherwise
    """
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        not_contains_patterns = expected.get('not_contains', [])
        if not not_contains_patterns:
            return 1.0
        pattern = '\n    '.join(not_contains_patterns)
        if pattern in content:
            return 0.0
        import re
        escaped_patterns = [re.escape(p.strip()) for p in not_contains_patterns]
        regex_pattern = '\\s*\\n\\s*'.join(escaped_patterns)
        if re.search(regex_pattern, content):
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_vscode_wordwrap__ca4beaed446890c3b888d4a40b939831(actual: str, expected: Dict, **options) -> float:
    """Check if VS Code settings.json has word wrap configured correctly.

    Args:
        actual: Path to settings.json file
        expected: Expected configuration rules (from expected.rules)
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel or actual_value != value:
            return 0.0
    return 1.0

def check_vscode_variable_rename__09b11b45(result, expected, **options):
    """
    Check if variable is renamed correctly throughout all functions in the file using AST parsing.

    This evaluator checks that the old variable name has been replaced with the new variable name
    across all functions in the file. It uses AST parsing to identify variable references and
    function parameters, ensuring accurate detection.

    Scoring:
    - 1.0: All occurrences of old_name have been renamed to new_name (no old_name remains)
    - 0.0-1.0: Partial credit based on the proportion of successful renames
    - 0.0: No renaming detected or no occurrences found
    """
    import ast
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        new_name = expected.get('new_name', '')
        old_name = expected.get('old_name', '')
        if not new_name or not old_name:
            return 0.0
        try:
            tree = ast.parse(content)
        except SyntaxError:
            return 0.0
        functions = [node for node in ast.walk(tree) if isinstance(node, ast.FunctionDef)]
        if not functions:
            return 0.0
        total_old_occurrences = 0
        total_new_occurrences = 0
        for func in functions:
            for node in ast.walk(func):
                if isinstance(node, ast.Name):
                    if node.id == old_name:
                        total_old_occurrences += 1
                    elif node.id == new_name:
                        total_new_occurrences += 1
                elif isinstance(node, ast.arg):
                    if node.arg == old_name:
                        total_old_occurrences += 1
                    elif node.arg == new_name:
                        total_new_occurrences += 1
        if total_old_occurrences == 0 and total_new_occurrences == 0:
            return 0.0
        if total_new_occurrences > 0 and total_old_occurrences == 0:
            return 1.0
        if total_old_occurrences > 0 and total_new_occurrences > 0:
            total_expected = total_old_occurrences + total_new_occurrences
            score = total_new_occurrences / total_expected
            return round(score, 2)
        if total_old_occurrences > 0 and total_new_occurrences == 0:
            return 0.0
        return 0.0
    except Exception as e:
        return 0.0

def check_workspace_folder_count__8a8ecf82(result: str, expected, **options) -> float:
    """
    Check if the workspace has the expected number of folders.

    Args:
        result (str): path to workspace file
        expected: dict containing "count" key with expected folder count
        **options: additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if not result:
        logger.warning('Result file path is None, returning 0.0')
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Error loading JSON from {result}: {e}')
        return 0.0
    folders = data.get('folders', [])
    actual_count = len(folders)
    expected_count = expected.get('count', 0)
    if actual_count == expected_count:
        return 1.0
    else:
        logger.debug(f'Folder count mismatch: expected {expected_count}, got {actual_count}')
        return 0.0

def check_vscode_settings_values__51833a76bce91777752996df0b22c33d(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pairs.

    Args:
        result: Actual settings dict from getter
        expected: Expected settings dict with 'settings' key containing key-value pairs to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on how many expected settings match
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_settings = expected.get('settings', {})
    if not expected_settings:
        return 0.0
    total = len(expected_settings)
    matched = 0
    for (key, expected_value) in expected_settings.items():
        if key in result and result[key] == expected_value:
            matched += 1
    return matched / total if total > 0 else 0.0

def check_vscode_line_content__da119c1e(result, expected, **options):
    """Check if a specific line contains expected content."""
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            lines = f.readlines()
        line_num = expected.get('line_num', 1)
        expected_pattern = expected.get('pattern', '')
        if line_num > len(lines):
            return 0.0
        actual_line = lines[line_num - 1].rstrip('\n')
        if expected_pattern in actual_line or actual_line.strip() == expected_pattern.strip():
            return 1.0
        return 0.0
    except Exception as e:
        return 0.0

def check_json_keybindings__c92c626f(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    data = direct_load_json(actual)
    if data is None or type(data) != list:
        data = skip_first_line_load_json(actual)
        if data is None or type(data) != list:
            return 0.0
    expected_binding = expected.get('expected')
    if not expected_binding or not isinstance(expected_binding, dict):
        return 0.0
    expected_key = expected_binding.get('key', '').lower()
    expected_command = expected_binding.get('command', '')
    if not expected_key or not expected_command:
        return 0.0
    for binding in data:
        if not isinstance(binding, dict):
            continue
        binding_key = binding.get('key', '').lower()
        binding_command = binding.get('command', '')
        if binding_key == expected_key and binding_command == expected_command:
            return 1.0
    return 0.0

def check_vscode_formatonsave__4df3f2e6(actual: str, rules: Dict, **options) -> float:
    """
    Check if VS Code format on save is enabled.

    Args:
        actual (str): Path to the settings.json file
        rules (Dict): Expected rules with 'expected_value' key

    Returns:
        float: 1.0 if format on save matches expected, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expected_value = rules.get('expected_value', False)
    actual_value = data.get('editor.formatOnSave', False)
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_vscode_function_name__75a1331b(result, expected, **options):
    """Check if function name is changed correctly throughout the file."""
    import re
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        expected_name = expected.get('function_name', '')
        new_func_pattern = 'def\\s+' + re.escape(expected_name) + '\\s*\\('
        has_new_function = re.search(new_func_pattern, content) is not None
        old_func_def_pattern = 'def\\s+bubble_sort\\s*\\('
        old_func_call_pattern = 'bubble_sort\\s*\\('
        has_old_func_def = re.search(old_func_def_pattern, content) is not None
        has_old_func_call = re.search(old_func_call_pattern, content) is not None
        if has_new_function and (not has_old_func_def) and (not has_old_func_call):
            return 1.0
        return 0.0
    except Exception as e:
        return 0.0

def check_vscode_indentation__06f73a877cbecd1e2a247622c3136810(result, expected, **options):
    """Compare file content with expected text.

    Args:
        result: Actual file content (string)
        expected: Expected content (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if content matches exactly, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_content = expected.get('content', '')
    if not expected_content:
        return 0.0
    if result == expected_content:
        return 1.0
    return 0.0

def check_vscode_setting__4269362e9e26c329da428a21ddabb7ea(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pair.

    Args:
        result: Dictionary from getter containing actual VS Code settings
        expected: Dictionary with 'key' and 'value' to check (from rules)
        **options: Additional options

    Returns:
        1.0 if the setting key has the expected value, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    try:
        if float(actual_value) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_vscode_contains_text__2fa9f229(result, expected, **options):
    """Check if file contains expected text strings."""
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        expected_strings = expected.get('contains', [])
        for s in expected_strings:
            if s not in content:
                return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_vscode_open_file__005f6ac9(actual: str, rules: dict, **options) -> float:
    """Check if the expected file is open in VSCode."""
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read().strip()
        expected = rules['expected']
        if expected in content:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking open file: {e}')
        return 0.0

def check_json_keybindings__c38e7539(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (from vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except:
        pass
    if data is None or not isinstance(data, list):
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            pass
    if data is None or not isinstance(data, list):
        return 0.0
    expected_binding = expected.get('expected')
    for binding in data:
        if isinstance(binding, dict):
            if all((binding.get(k) == v for (k, v) in expected_binding.items())):
                return 1.0
    return 0.0

def check_json_keybindings__3742c48e(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if not isinstance(data, list):
            data = None
    except:
        pass
    if data is None:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
            if not isinstance(data, list):
                data = None
        except:
            pass
    if data is None:
        return 0.0
    if 'expected' in expected:
        expected_binding = expected['expected']
    else:
        expected_binding = expected
    expected_key = expected_binding.get('key')
    expected_command = expected_binding.get('command')
    expected_when = expected_binding.get('when')
    if expected_key is None or expected_command is None or expected_when is None:
        return 0.0
    for binding in data:
        if not isinstance(binding, dict):
            continue
        if binding.get('key') == expected_key and binding.get('command') == expected_command and (binding.get('when') == expected_when):
            return 1.0
    return 0.0

def check_vscode_wordwrap__ff65f872(actual: str, rules: Dict, **options) -> float:
    """
    Check if VS Code word wrap is enabled.

    Args:
        actual (str): Path to the settings.json file
        rules (Dict): Expected rules with 'expected_value' key

    Returns:
        float: 1.0 if word wrap matches expected, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expected_value = rules.get('expected_value', '')
    actual_value = data.get('editor.wordWrap', '')
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_vscode_setting__37a81987aef2f239d369b6635778b9ec(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pair.

    Args:
        result: Dictionary from getter containing actual VS Code settings
        expected: Dictionary with 'key' and 'value' to check (from rules)
        **options: Additional options

    Returns:
        1.0 if the setting key has the expected value, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    try:
        if float(actual_value) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_vscode_tabsize__a6b4939a2d71540bebc75ede5d00109a(actual: str, expected: Dict, **options) -> float:
    """Check if VS Code settings.json has tab size configured correctly.

    Args:
        actual: Path to settings.json file
        expected: Expected configuration rules (from expected.rules)
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel or actual_value != value:
            return 0.0
    return 1.0

def check_vscode_config_files__17da0621(actual: str, rules: dict, **options) -> float:
    """Check if workspace has config files."""
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read().strip()
        expected_file = rules.get('expected_file', 'settings.json')
        if expected_file in content:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking config files: {e}')
        return 0.0

def check_workspace_folder_order__40d497fa35fd103d8064472e794b77f6(actual: str, expected: Dict, **options) -> float:
    """
    Check if workspace file contains expected folders in exact order.

    Args:
        actual (str): path to workspace file
        expected (Dict): expected configuration with 'folders' key (ordered)

    Returns:
        float: 1.0 if folders match in exact order, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expected_folders = expected.get('folders', [])
    actual_folders = data.get('folders', [])
    if len(expected_folders) != len(actual_folders):
        return 0.0
    for (exp_folder, act_folder) in zip(expected_folders, actual_folders):
        if exp_folder.get('path') != act_folder.get('path'):
            return 0.0
    return 1.0

def check_json_keybindings__8bf59803(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 4: Remove Ctrl+P shortcut for workbench.action.quickOpen

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if not actual or not isinstance(actual, str):
        return 0.0
    if not os.path.exists(actual):
        return 0.0
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            return 0.0
    if data is None or not isinstance(data, list):
        return 0.0
    expected_entry = expected['expected']
    expected_key = expected_entry.get('key')
    expected_command = expected_entry.get('command')
    if expected_command and (not expected_command.startswith('-')):
        return 0.0
    for entry in data:
        if isinstance(entry, dict):
            if entry.get('key') == expected_key and entry.get('command') == expected_command:
                return 1.0
    return 0.0

def check_workspace_folders__328ca95134872b7162fe90bdb76e5fde(actual: str, expected: Dict, **options) -> float:
    """
    Check if workspace file contains expected folders.

    Args:
        actual (str): path to workspace file
        expected (Dict): expected configuration with 'folders' key

    Returns:
        float: 1.0 if folders match, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expected_folders = expected.get('folders', [])
    actual_folders = data.get('folders', [])
    if len(expected_folders) != len(actual_folders):
        return 0.0
    expected_paths = {folder['path'] for folder in expected_folders}
    actual_paths = {folder['path'] for folder in actual_folders}
    if expected_paths == actual_paths:
        return 1.0
    return 0.0

def check_vscode_open_file__149df8b642d00f3a86d140b4582f20cc(actual: str, expected: Dict, **options) -> float:
    """Check if the expected file is currently open in VSCode.

    Args:
        actual: Path to the file containing the currently open file info
        expected: Dict with 'expected' key containing expected filename
        **options: Additional options

    Returns:
        float: 1.0 if expected file is open, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
        expected_file = expected.get('expected', '')
        if expected_file in actual_text:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking open file: {e}')
        return 0.0

def check_json_keybindings__9f7d30da(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 6: Remove Ctrl+Shift+F shortcut for workbench.action.findInFiles

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            return 0.0
    except:
        pass
    try:
        with open(actual, 'r') as f:
            f.readline()
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            return 0.0
    except:
        pass
    return 0.0

def check_keybinding_editor_to_terminal__d498cab59d97f4e17938744aa091cf0b(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding is configured in VSCode keybindings.json.
    This variation checks for editor-to-terminal focus switch.

    Args:
        actual (str): path to result keybindings.json file
        expected (Dict): expected dict with 'expected' key containing the keybinding config

    Return:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if actual is None:
        return 0.0
    if not isinstance(actual, str):
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    else:
        return 0.0
    expected_binding = expected['expected']
    if expected_binding in data:
        return 1.0
    else:
        return 0.0

def check_multiple_workspaces__f54bf781c32dea93b80d6edc50522687(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if multiple workspace files exist and are valid.

    Args:
        result: Dict from getter with 'workspaces', 'total_count', 'exists_count', 'valid_count' keys
        expected: Dict with 'expected_count' and 'expected_paths' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_count = expected.get('expected_count', 0)
    expected_paths = expected.get('expected_paths', [])
    exists_count = result.get('exists_count', 0)
    valid_count = result.get('valid_count', 0)
    workspaces = result.get('workspaces', [])
    score = 0.0
    if exists_count >= expected_count:
        score += 0.5
    if valid_count >= expected_count:
        score += 0.3
    if expected_paths:
        existing_paths = {w['path'] for w in workspaces if w['exists']}
        expected_paths_set = set(expected_paths)
        if expected_paths_set.issubset(existing_paths):
            score += 0.2
    return score

def check_json_keybindings__ec7a7a17(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if not actual or not isinstance(actual, str):
        return 0.0
    if not os.path.exists(actual):
        return 0.0
    file_path = actual

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    keybindings_data = None
    for func in [direct_load_json, skip_first_line_load_json]:
        keybindings_data = func(file_path)
        if keybindings_data is not None and isinstance(keybindings_data, list):
            break
    if keybindings_data is None:
        return 0.0
    expected_binding = expected.get('expected')
    if not expected_binding:
        return 0.0
    expected_key = expected_binding.get('key')
    expected_command = expected_binding.get('command')
    for binding in keybindings_data:
        if not isinstance(binding, dict):
            continue
        if binding.get('key') == expected_key and binding.get('command') == expected_command:
            return 1.0
    return 0.0

def check_json_keybindings__45d3929c(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 2: Remove Ctrl+/ shortcut for editor.action.commentLine

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except:
        pass
    if data is None or not isinstance(data, list):
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            pass
    if data is None or not isinstance(data, list):
        return 0.0
    expected_entry = expected['expected']
    if expected_entry in data:
        return 1.0
    else:
        return 0.0

def check_workspace_has_folder__33d8878a(result: str, expected, **options) -> float:
    """
    Check if the workspace contains a specific folder by path.

    Args:
        result (str): path to workspace file
        expected: dict containing "folder_path" key with the folder path to check
        **options: additional options

    Returns:
        float: 1.0 if folder exists in workspace, 0.0 otherwise
    """
    if not result:
        logger.warning('Result file path is None, returning 0.0')
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Error loading JSON from {result}: {e}')
        return 0.0
    folders = data.get('folders', [])
    expected_path = expected.get('folder_path', '')
    for folder in folders:
        folder_path = folder.get('path', '')
        normalized_folder = os.path.normpath(folder_path)
        normalized_expected = os.path.normpath(expected_path)
        if normalized_folder == normalized_expected:
            return 1.0
        if os.path.basename(normalized_folder) == normalized_expected:
            return 1.0
        if normalized_folder.endswith(os.sep + normalized_expected):
            return 1.0
        if normalized_folder.endswith(normalized_expected):
            return 1.0
    logger.debug(f"Folder with path '{expected_path}' not found in workspace")
    return 0.0

def check_vscode_indentation__219e3148d77134036e282c6dd4b41d12(result, expected, **options):
    """Compare file content with expected text.

    Args:
        result: Actual file content (string)
        expected: Expected content (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if content matches exactly, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_content = expected.get('content', '')
    if not expected_content:
        return 0.0
    if result == expected_content:
        return 1.0
    return 0.0

def check_json_keybindings__a807c39e(vm_file: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 7: Disable Ctrl+K Ctrl+C shortcut for editor.action.addCommentLine

    Args:
        vm_file (str): path to keybindings.json file from vm_file result
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    file_path = vm_file
    if not file_path or not isinstance(file_path, str):
        return 0.0
    data = None
    try:
        with open(file_path, 'r') as f:
            content = f.read()
            data = json.loads(content)
        if data is not None and type(data) == list:
            pass
        else:
            data = None
    except:
        data = None
    if data is None:
        try:
            with open(file_path, 'r') as f:
                f.readline()
                content = f.read()
                data = json.loads(content)
            if data is None or type(data) != list:
                data = None
        except:
            data = None
    if data is None:
        return 0.0
    expected_entry = expected['expected']
    if expected_entry in data:
        return 1.0
    else:
        return 0.0

def check_workspace_folders__ff3abf45ba7fc3112a6dd2bb6da6358b(actual: str, rules: Dict, **options) -> float:
    """
    Check if VSCode workspace contains the required folders.

    This function flexibly checks for folder membership:
    - Handles both absolute and relative paths
    - Order-independent comparison
    - Only checks 'path' property, ignoring additional properties like 'name'
    - Normalizes paths to handle /home/user/folder vs folder

    Args:
        actual (str): path to the workspace JSON file
        rules (Dict): dict containing 'expected' with the required folders list

    Returns:
        float: 1.0 if all required folders are present, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    actual_folders = data.get('folders', [])
    if not isinstance(actual_folders, list):
        return 0.0
    expected = rules.get('expected', {})
    expected_folders = expected.get('folders', [])
    if not isinstance(expected_folders, list):
        return 0.0
    actual_paths = set()
    for folder in actual_folders:
        if isinstance(folder, dict) and 'path' in folder:
            path = folder['path']
            normalized = path.strip()
            if normalized.startswith('/home/user/'):
                normalized = normalized.replace('/home/user/', '', 1)
            actual_paths.add(normalized)
    expected_paths = set()
    for folder in expected_folders:
        if isinstance(folder, dict) and 'path' in folder:
            path = folder['path']
            normalized = path.strip()
            if normalized.startswith('/home/user/'):
                normalized = normalized.replace('/home/user/', '', 1)
            expected_paths.add(normalized)
    if expected_paths.issubset(actual_paths):
        return 1.0
    return 0.0

def check_json_keybindings__c2ab0a78(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 1: Remove Ctrl+Shift+K shortcut for editor.action.deleteLines

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            else:
                return 0.0
    except:
        pass
    try:
        with open(actual, 'r') as f:
            f.readline()
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            else:
                return 0.0
    except:
        pass
    return 0.0

def check_json_keybindings__093ec527(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 3: Disable Ctrl+B shortcut for workbench.action.toggleSidebarVisibility

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            else:
                return 0.0
    except:
        pass
    try:
        with open(actual, 'r') as f:
            f.readline()
            data = json.load(f)
        if isinstance(data, list):
            expected_entry = expected['expected']
            if expected_entry in data:
                return 1.0
            else:
                return 0.0
    except:
        pass
    return 0.0

def check_json_keybindings__940d9534(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except:
            return 0.0
    if data is None or not isinstance(data, list):
        return 0.0
    expected_binding = expected.get('expected')
    if not expected_binding:
        return 0.0
    expected_key = expected_binding.get('key')
    expected_command = expected_binding.get('command')
    if not expected_key or not expected_command:
        return 0.0
    for binding in data:
        if isinstance(binding, dict):
            binding_key = binding.get('key', '').lower()
            binding_command = binding.get('command', '')
            if binding_key == expected_key.lower() and binding_command == expected_command:
                return 1.0
    return 0.0

def check_vscode_setting__a46d28041c0fe91cc8b08a2032679781(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain expected key-value pair.

    Args:
        result: Dictionary from getter containing actual VS Code settings
        expected: Dictionary with 'key' and 'value' to check (from rules)
        **options: Additional options

    Returns:
        1.0 if the setting key has the expected value, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    try:
        if float(actual_value) == float(expected_value):
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_vscode_markdown_files__77849053(actual: str, rules: dict, **options) -> float:
    """Check if workspace has markdown files."""
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read().strip()
        expected_file = rules.get('expected_file', 'README.md')
        if expected_file in content:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking markdown files: {e}')
        return 0.0

def check_vscode_workspace__6efce4eb1d6563b7b443059d330de6aa(actual: str, rules: Dict, **options) -> float:
    """Check if the expected workspace folder is open in VSCode.

    Args:
        actual: Path to the file containing the workspace info
        rules: Dict with 'expected' key containing expected workspace folder name
        **options: Additional options

    Returns:
        float: 1.0 if expected workspace is open, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
        expected = rules.get('expected', '')
        if expected in actual_text:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking workspace: {e}')
        return 0.0

def check_workspace_folders__9f48007a04e5b61abc6f8fde0b255a1d(actual: str, expected: Dict, **options) -> float:
    """
    Check if workspace file contains expected folders.

    Args:
        actual (str): path to workspace file
        expected (Dict): expected configuration with 'folders' key

    Returns:
        float: 1.0 if folders match, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expected_folders = expected.get('folders', [])
    actual_folders = data.get('folders', [])
    if len(expected_folders) != len(actual_folders):
        return 0.0
    expected_paths = {folder['path'] for folder in expected_folders}
    actual_paths = {folder['path'] for folder in actual_folders}
    if expected_paths == actual_paths:
        return 1.0
    return 0.0

def check_json_keybindings__76d7a1f4(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file (vm_file result)
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    file_path = str(actual)
    file_path = os.path.abspath(file_path)
    if not os.path.exists(file_path):
        return 0.0
    if not os.path.isfile(file_path):
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    data = None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(file_path)
        if data is not None and type(data) == list:
            break
    if data is None or type(data) != list:
        return 0.0
    expected_binding = expected.get('expected')
    for item in data:
        if isinstance(item, dict):
            if all((expected_binding.get(k) == item.get(k) for k in expected_binding.keys())):
                return 1.0
    return 0.0

def check_vscode_autosave__aec709e9(actual: str, rules: Dict, **options) -> float:
    """
    Check if VS Code auto-save is enabled.

    Args:
        actual (str): Path to the settings.json file
        rules (Dict): Expected rules with 'expected_value' key

    Returns:
        float: 1.0 if auto-save is enabled, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expected_value = rules.get('expected_value', '')
    actual_value = data.get('files.autoSave', '')
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_vscode_format_on_paste__6ee0182d(actual: str, rules: dict, **options) -> float:
    """
    Check if VSCode settings.json has editor.formatOnPaste set to the expected value.

    Args:
        actual (str): path to settings.json file
        rules (dict): expected configuration rules with format_on_paste

    Returns:
        float: score between 0.0 and 1.0
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Failed to load settings.json: {e}')
        return 0.0
    expected_value = rules.get('format_on_paste', True)
    score = 0.0
    if 'editor.formatOnPaste' in data:
        if data['editor.formatOnPaste'] == expected_value:
            score = 1.0
    return score

def check_vscode_open_file__c60120ed87e101802d8b786344162dfa(actual: str, rules: Dict, **options) -> float:
    """Check if the expected file is currently open in VSCode.

    Args:
        actual: Path to the file containing the currently open file info
        rules: Dict with 'expected' key containing expected filename
        **options: Additional options

    Returns:
        float: 1.0 if expected file is open, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
        expected = rules.get('expected', '')
        if not expected:
            return 0.0
        if actual_text.endswith(expected) or actual_text.endswith('/' + expected):
            if '/home/user/project/' in actual_text or actual_text == expected:
                return 1.0
        expected_full_path = f'/home/user/project/{expected}'
        if actual_text == expected_full_path:
            return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking open file: {e}')
        return 0.0

def check_vscode_workspace_files__4ab817c3(actual: str, rules: dict, **options) -> float:
    """Check if workspace contains expected files."""
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        expected_files = rules.get('expected_files', [])
        score = 0.0
        for expected_file in expected_files:
            if expected_file in content:
                score += 1.0 / len(expected_files)
        return score
    except Exception as e:
        logger.error(f'Error checking workspace files: {e}')
        return 0.0

def check_vscode_setting__b07969a1(actual: str, expected: Dict, **options) -> float:
    """Check if a specific setting has the expected value in VS Code.

    Args:
        actual (str): path to result JSON file containing settings
        expected (Dict): expected dict with 'setting_key' and 'setting_value' keys
        **options: Additional options

    Return:
        float: the score (1.0 if setting is correct, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    setting_key = expected.get('setting_key', '')
    setting_value = expected.get('setting_value', '')
    if settings.get(setting_key) == setting_value:
        return 1.0
    return 0.0

def check_vscode_line_content__9d7d44f7(result, expected, **options):
    """Check if a specific line contains expected content."""
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            lines = f.readlines()
        line_num = expected.get('line_num', 1)
        expected_pattern = expected.get('pattern', '')
        if line_num > len(lines):
            return 0.0
        actual_line = lines[line_num - 1].rstrip('\n')
        if actual_line.strip() == expected_pattern.strip():
            return 1.0
        return 0.0
    except Exception as e:
        return 0.0

def check_json_keybindings__e8fad444(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding exists in VSCode keybindings.json.
    Variation 9: Disable Ctrl+H shortcut for editor.action.startFindReplaceAction

    Args:
        actual (str): path to keybindings.json file
        expected (Dict): expected keybinding entry with 'expected' key

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if not actual:
        return 0.0
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
        if data is not None and type(data) == list:
            pass
        else:
            data = None
    except:
        data = None
    if data is None:
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
            if data is not None and type(data) == list:
                pass
            else:
                data = None
        except:
            data = None
    if data is None or type(data) != list:
        return 0.0
    expected_entry = expected['expected']
    if expected_entry in data:
        return 1.0
    else:
        return 0.0

def check_vscode_git_exclude__3826e4daac10a542e57e7cc4a2b6b258(actual: str, expected: Dict, **options) -> float:
    """Check if VS Code settings.json has .git folder exclusion configured.

    Args:
        actual: Path to settings.json file
        expected: Expected configuration rules (from expected.rules)
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expect = expected.get('expected', {})
    for (key, value) in expect.items():
        sentinel = object()
        actual_value = data.get(key, sentinel)
        if actual_value is sentinel:
            return 0.0
        if isinstance(value, dict):
            if not isinstance(actual_value, dict):
                return 0.0
            for (k, v) in value.items():
                if actual_value.get(k) != v:
                    return 0.0
        elif actual_value != value:
            return 0.0
    return 1.0

def check_vscode_file_count__7c2e5da9(actual: str, rules: dict, **options) -> float:
    """Check if workspace has the expected count of file types."""
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read().strip()
        file_type = rules.get('file_type', '.py')
        expected_count = rules.get('expected_count', 1)
        count = content.count(file_type)
        if count >= expected_count:
            return 1.0
        else:
            return count / expected_count
    except Exception as e:
        logger.error(f'Error checking file count: {e}')
        return 0.0

def check_vscode_tabsize__44a94327(actual: str, rules: Dict, **options) -> float:
    """
    Check if VS Code tab size is set to the expected value.

    Args:
        actual (str): Path to the settings.json file
        rules (Dict): Expected rules with 'expected_size' key

    Returns:
        float: 1.0 if tab size matches, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception:
        return 0.0
    expected_size = rules.get('expected_size', 0)
    actual_size = data.get('editor.tabSize', 0)
    if actual_size == expected_size:
        return 1.0
    return 0.0

def check_vscode_indentation__57096d124fd446ca6cba7d1316aac1bd(result, expected, **options):
    """Compare file content with expected text.

    Args:
        result: Actual file content (string)
        expected: Expected content (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if content matches exactly, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_content = expected.get('content', '')
    if not expected_content:
        return 0.0
    if result == expected_content:
        return 1.0
    return 0.0

def check_vscode_open_file__e925f44c(actual: str, expected: Dict, **options) -> float:
    """Check if the correct file is open in VS Code.

    Args:
        actual (str): path to result text file containing the open file name
        expected (Dict): expected dict with 'file_path' key
        **options: Additional options

    Return:
        float: the score (1.0 if correct file is open, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
    except Exception:
        return 0.0
    expected_file = expected.get('file_path', '')
    if expected_file in actual_text:
        return 1.0
    return 0.0

def check_keybinding_sidebar_toggle__9bbee30a9c2c82b51b389394e0233d5e(actual: str, expected: Dict, **options) -> float:
    """
    Check if a specific keybinding is configured in VSCode keybindings.json.
    This variation checks for sidebar toggle keybinding.

    Args:
        actual (str): path to result keybindings.json file
        expected (Dict): expected dict with 'expected' key containing the keybinding config

    Return:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if actual is None:
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    else:
        return 0.0
    expected_binding = expected['expected']
    for binding in data:
        if not isinstance(binding, dict):
            continue
        key_match = binding.get('key') == expected_binding.get('key')
        command_match = binding.get('command') == expected_binding.get('command')
        when_expected = expected_binding.get('when')
        when_actual = binding.get('when')
        if when_expected is not None:
            when_match = when_actual == when_expected
        else:
            when_match = True
        if key_match and command_match and when_match:
            return 1.0
    return 0.0

def check_json_keybindings__9c5a8886(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific keybinding is present in VS Code keybindings.json

    Args:
        actual: path to keybindings.json file
        expected: dict containing the expected keybinding (from rules)

    Returns:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if actual is None:
        return 0.0
    if not os.path.exists(actual):
        return 0.0
    data = None
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except (json.JSONDecodeError, IOError, FileNotFoundError):
        try:
            with open(actual, 'r') as f:
                f.readline()
                data = json.load(f)
        except (json.JSONDecodeError, IOError, FileNotFoundError):
            return 0.0
    if data is None or not isinstance(data, list):
        return 0.0
    expected_binding = expected.get('expected')
    if expected_binding in data:
        return 1.0
    else:
        return 0.0

def check_vscode_settings_autosave__06d1442bba585f1fb8ea50a4956df674(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VS Code settings contain the expected autosave configuration.

    Args:
        result: Dictionary containing actual settings.json content
        expected: Dictionary with 'key' and 'value' to check
        **options: Additional options

    Returns:
        float: 1.0 if settings match expected, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_key = expected.get('key')
    expected_value = expected.get('value')
    if expected_key is None:
        logger.error('Expected key not specified')
        return 0.0
    actual_value = result.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_workspace_folders(actual, expected, **options) -> float:
    """
    Check if workspace file contains all required folders with flexible path matching.

    This function handles VSCode workspace files which can store paths in multiple formats:
    - Absolute paths: /home/user/project
    - Relative paths with ./: ./project
    - Relative paths without prefix: project

    Args:
        actual (str): path to workspace file
        expected (dict): expected configuration with 'expected' key containing folder structure
        **options: additional options

    Returns:
        float: 1.0 if all expected folders are present, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        return 0.0
    expect = expected.get('expected', {})
    expected_folders = expect.get('folders', [])
    if not expected_folders:
        return 0.0
    actual_folders = data.get('folders', [])
    if not actual_folders:
        return 0.0
    actual_paths = set()
    for folder in actual_folders:
        if isinstance(folder, dict) and 'path' in folder:
            path = folder['path']
            normalized = normalize_path(path, actual)
            actual_paths.add(normalized)
    expected_paths = set()
    for folder in expected_folders:
        if isinstance(folder, dict) and 'path' in folder:
            path = folder['path']
            normalized = normalize_path(path, actual)
            expected_paths.add(normalized)
    if expected_paths.issubset(actual_paths):
        return 1.0
    return 0.0

def check_keybinding_terminal_toggle__326285452dea673d67652e18309c6e13(actual: Union[str, bytes, None], expected: Dict, **options) -> float:
    """
    Check if a specific keybinding is configured in VSCode keybindings.json.
    This variation checks for terminal toggle keybinding.

    Args:
        actual (Union[str, bytes, None]): path to result keybindings.json file or file content
        expected (Dict): expected dict with 'expected' key containing the keybinding config

    Return:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if actual is None:
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except:
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except:
            return None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    else:
        return 0.0
    if data is None:
        return 0.0
    expected_binding = expected['expected']
    for binding in data:
        if binding.get('key') == expected_binding.get('key') and binding.get('command') == expected_binding.get('command') and (binding.get('when') == expected_binding.get('when')):
            return 1.0
    return 0.0

def check_workspace_excludes_folder__2068a9c3(result: str, expected, **options) -> float:
    """
    Check if the workspace does NOT contain a specific folder by path.

    Args:
        result (str): path to workspace file
        expected: dict containing "folder_path" key with the folder path that should NOT be present
        **options: additional options

    Returns:
        float: 1.0 if folder does NOT exist in workspace, 0.0 if it exists
    """
    if not result:
        logger.warning('Result file path is None, returning 0.0')
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            data = json.load(f)
    except Exception as e:
        logger.error(f'Error loading JSON from {result}: {e}')
        return 0.0
    folders = data.get('folders', [])
    excluded_path = expected.get('folder_path', '')
    for folder in folders:
        if folder.get('path') == excluded_path:
            logger.debug(f"Folder with path '{excluded_path}' found in workspace (should be excluded)")
            return 0.0
    return 1.0

def check_vscode_contains_text__39395653(result, expected, **options):
    """Check if file contains a proper docstring after function definition."""
    import ast
    import re
    if not result:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read()
        try:
            tree = ast.parse(content)
            for node in ast.walk(tree):
                if isinstance(node, ast.FunctionDef):
                    docstring = ast.get_docstring(node)
                    if docstring:
                        if 'bubble' in docstring.lower():
                            return 1.0
        except SyntaxError:
            pass
        pattern = 'def\\s+\\w+\\s*\\([^)]*\\)\\s*:\\s*\\n\\s+"""([^"]|"(?!""))*bubble([^"]|"(?!""))*"""'
        if re.search(pattern, content, re.IGNORECASE | re.DOTALL):
            return 1.0
        pattern_single = 'def\\s+\\w+\\s*\\([^)]*\\)\\s*:\\s*\\n\\s+"""[^"]*bubble[^"]*"""'
        if re.search(pattern_single, content, re.IGNORECASE):
            return 1.0
        return 0.0
    except Exception as e:
        return 0.0

def check_vscode_setting__232783f17230c635388a2d96b7097023(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if a specific VSCode setting has the expected value.

    Args:
        result: Dict containing the actual settings
        expected: Dict with 'key' and 'value' fields specifying what to check

    Returns:
        1.0 if the setting matches, 0.0 otherwise
    """
    if not result:
        return 0.0
    key = expected.get('key')
    expected_value = expected.get('value')
    if key not in result:
        return 0.0
    actual_value = result.get(key)
    if actual_value == expected_value:
        return 1.0
    return 0.0

def check_keybinding_reopen_editor__252b583b8604f0bd370b6e7ab984580e(actual: Union[str, None], expected: Dict, **options) -> float:
    """
    Check if a specific keybinding is configured in VSCode keybindings.json.
    This variation checks for reopening closed editor keybinding.

    Args:
        actual (Union[str, None]): path to result keybindings.json file (can be None if fetch fails)
        expected (Dict): expected dict with 'expected' key containing the keybinding config

    Return:
        float: 1.0 if keybinding exists, 0.0 otherwise
    """
    if actual is None:
        logging.warning('File fetch failed: actual is None')
        return 0.0

    def direct_load_json(fp):
        try:
            with open(fp, 'r') as f:
                data = json.load(f)
            return data
        except Exception as e:
            logging.debug(f'direct_load_json failed: {e}')
            return None

    def skip_first_line_load_json(fp):
        try:
            with open(fp, 'r') as f:
                f.readline()
                data = json.load(f)
            return data
        except Exception as e:
            logging.debug(f'skip_first_line_load_json failed: {e}')
            return None
    for func in [direct_load_json, skip_first_line_load_json]:
        data = func(actual)
        if data is not None and type(data) == list:
            break
    else:
        logging.error(f'Failed to parse JSON from file: {actual}')
        return 0.0
    expected_binding = expected['expected']
    if any((kb == expected_binding for kb in data)):
        return 1.0
    else:
        return 0.0

def check_vscode_format_on_save__7a482b52(actual: str, expected: Dict, **options) -> float:
    """Check if format on save is disabled in VS Code settings.

    Args:
        actual (str): path to result JSON file containing settings
        expected (Dict): expected dict with 'format_on_save' key (boolean)
        **options: Additional options

    Return:
        float: the score (1.0 if setting matches, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            settings = json.load(f)
    except Exception:
        return 0.0
    expected_value = expected.get('format_on_save', False)
    actual_value = settings.get('editor.formatOnSave', True)
    if actual_value == expected_value:
        return 1.0
    return 0.0
