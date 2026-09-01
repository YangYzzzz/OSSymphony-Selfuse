"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_thunderbird_filter__f57034f8', 'check_thunderbird_digest_subjects__5a1ab509', 'check_thunderbird_imap_settings__bd92e86076085c3fd0cc432ab151ca8f', 'check_thunderbird_subject_date__2284321c', 'check_tb_filtered_emails__01abb20597fbc4723613687c7c60c0d3', 'check_thunderbird_three_cols__07170c7f', 'check_thunderbird_txt_files__dd834e2c641c360491704c9db4c6f362', 'check_thunderbird_server_check__734b4908962171d93d8cfb6878a651d2', 'check_thunderbird_timeout__0d04b8a7', 'check_thunderbird_email_deleted_and_folder_created__3c15b6d0', 'check_thunderbird_date_from_subject__a92c1922', 'check_thunderbird_check_new_mail__625e2b37', 'check_thunderbird_filter__2de049a0', 'check_email_domains__8d2fdf6f2182283d571d666af9c9e85f', 'check_thunderbird_workflow_and_amazon_tab__0f025d55a86976e1ae415a210c74f719', 'check_thunderbird_oldest_three__0c7bac1c', 'check_email_count__fcf6b1e44de70adba5822f1fa1e262b7', 'check_contact_emails__0fa83097089d9008d43d2c029fbbd194', 'check_thunderbird_filter__9b7bc335', 'check_thunderbird_prefs__249498a7cd34d43ebd7ef238b7ba6bee', 'check_thunderbird_drafts_picker__17405a90', 'check_thunderbird_filter__ce00a590', 'check_thunderbird_filter__07f67672', 'check_thunderbird_folder_exists_and_email_starred__7e32e736', 'check_thunderbird_cc_emails__f218f3c2', 'check_thunderbird_sender_count__44f17dcf', 'check_all_emails_starred__d14b39d6', 'check_thunderbird_filter_created__8a2db6ad', 'check_thunderbird_amazon_deleted__4bba48e3', 'check_thunderbird_folder_and_empty_bills__a91026e2', 'check_first_email__105b8e17', 'check_thunderbird_compact_folders__e06bf170da6e0042eda0ceb7f1f9f833', 'check_thunderbird_subject__c9ce3f52', 'check_tb_digest__b6ecd864d1dd0f24d97ae590f2c596e8', 'check_tb_summary__f5f661c56b53d0549f4c6bd9201ba899', 'check_email_column__3631675f', 'check_email_subjects_list__df60c313c1f3b2e712e9756c47386ddc', 'check_all_emails_read__20eb9ee4', 'check_tb_sender_counts__6a3e0f1dfc20fc63f235dcd50f4dfaaf', 'check_thunderbird_trash_folder__550830d3', 'check_thunderbird_emails__f5c13cdd', 'check_thunderbird_filter__793e1ffa', 'check_tb_first3_emails__843811624c18ed6ed65e7d4d877f3122', 'check_thunderbird_all_addresses__d8fe40b3', 'check_thunderbird_pdfs_uploaded_to_gdrive__69a89619e7273cf1d986af62cad42166', 'check_thunderbird_filter__e2c8667a', 'check_thunderbird_prefs__5b041e31e4566087bb3bbe0b347cd7ff', 'check_thunderbird_reply_settings__616461c5ba008feaf990bbd287063304', 'check_tb_weekly__264d715cc53b32b5aee947baaf88b841', 'check_name_email_table__6e2aeb6fcb42c45735613c9af17f7e4d', 'check_email_not_in_recipients__c9ce3f52', 'check_thunderbird_filter__947ae516', 'check_thunderbird_unique_senders__0399616e', 'check_thunderbird_prefs__1f5ab547a29208fe81de884f80a1716d', 'check_thunderbird_full_name__cd29a448', 'check_thunderbird_reply_on_top__f50e981c', 'check_email_count_file__c95564e31fba18e60360c502646bdd5f', 'check_tb_oldest3__dcd1e9d2ff68b2e8e650947cdbb3db16', 'check_thunderbird_prefs__f0f94c7477cc672daff3cf5ec91ccac4', 'check_test_emails_exported__a9ff16fc1be9331aa02283bf7c168b3e', 'check_thunderbird_two_folders__386310fb', 'check_thunderbird_attachment_complete__d38192b0', 'check_thunderbird_filter__07c1111d', 'check_thunderbird_filter__782e51e3', 'check_thunderbird_filter__1da45b53', 'check_thunderbird_domain_count__1744c07b', 'check_thunderbird_filter__af973420', 'check_tb_email_summary__095c028a04cb1496af16d82dbb5f9c21', 'check_tb_enotices__5ef2369a29068457d8de24f5079f30bb', 'check_thunderbird_account_name__dced30f9', 'check_thunderbird_prefs__4d8c063289348a4e74b61401db33c326', 'check_thunderbird_all_emails_read__a63fb94e', 'check_thunderbird_draft_attachment__d38192b0_aug9', 'check_thunderbird_socket_timeout__a3ecd62701cd6d8575ca51c05b0fdb89', 'check_thunderbird_bcc__497f2c49']

def check_thunderbird_filter__f57034f8(result: str, rules: Dict[str, any]) -> float:
    """
    Check Thunderbird filter configuration with support for partial matching.

    Args:
        result (str): Path to filter definition file (msgFilterRules.dat)
        rules (Dict[str, any]): Expected filter properties including:
            - name: Filter name (exact match)
            - enabled: "yes" or "no" (exact match)
            - action: Action type (exact match, e.g., "Move to folder")
            - actionValue_contains: String that should be in actionValue (partial match)
            - condition_contains: String that should be in at least one condition (partial match)

    Returns:
        float: 1.0 if all requirements match, 0.0 otherwise
    """
    if result is None:
        logger.warning('Filter file result is None')
        return 0.0
    filters = []
    try:
        with open(result) as f:
            current_filter = None
            for line in f:
                line = line.strip()
                if line.startswith('name='):
                    if current_filter is not None:
                        filters.append(current_filter)
                    current_filter = {}
                    current_filter['name'] = _value_processor(line[5:].strip('"'))
                elif current_filter is not None:
                    if line.startswith('enabled='):
                        current_filter['enabled'] = _value_processor(line[8:].strip('"'))
                    elif line.startswith('type='):
                        current_filter['type'] = _value_processor(line[5:].strip('"'))
                    elif line.startswith('action='):
                        current_filter['action'] = _value_processor(line[7:].strip('"'))
                    elif line.startswith('actionValue='):
                        current_filter['actionValue'] = _value_processor(line[12:].strip('"'))
                    elif line.startswith('condition='):
                        condition_str = _value_processor(line[10:].strip('"'))
                        logger.debug('FILTER CONDITION: %s', condition_str)
                        conditions = _condition_pattern.findall(condition_str)
                        logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                        current_filter['condition'] = conditions
            if current_filter is not None:
                filters.append(current_filter)
    except Exception as e:
        logger.error('Error parsing filter file: %s', e)
        return 0.0
    logger.info('Parsed filters: %s', filters)
    for flt in filters:
        logger.debug('Checking filter: %s', flt)
        if 'name' in rules and flt.get('name') != rules['name']:
            logger.debug('Name mismatch: expected %s, got %s', rules['name'], flt.get('name'))
            continue
        if 'enabled' in rules and flt.get('enabled') != rules['enabled']:
            logger.debug('Enabled mismatch: expected %s, got %s', rules['enabled'], flt.get('enabled'))
            continue
        if 'action' in rules and flt.get('action') != rules['action']:
            logger.debug('Action mismatch: expected %s, got %s', rules['action'], flt.get('action'))
            continue
        if 'actionValue_contains' in rules:
            action_value = flt.get('actionValue', '')
            expected_substring = rules['actionValue_contains']
            if expected_substring not in action_value:
                logger.debug("ActionValue doesn't contain '%s': %s", expected_substring, action_value)
                continue
        if 'condition_contains' in rules:
            conditions = flt.get('condition', [])
            expected_substring = rules['condition_contains']
            if not any((expected_substring in cond for cond in conditions)):
                logger.debug("No condition contains '%s': %s", expected_substring, conditions)
                continue
        logger.info('Filter matches all requirements: %s', flt)
        return 1.0
    logger.warning('No filter matched the requirements')
    return 0.0

def check_thunderbird_digest_subjects__5a1ab509(result, expected, **options):
    """Compare digest subject list against expected values.

    Checks that all expected subjects are present in the result in chronological order.
    Allows for partial credit if some subjects are found.

    Args:
        result: List of subject lines from getter
        expected: Dict with 'expected_subjects' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_subjects = expected.get('expected_subjects', [])
    if not expected_subjects:
        return 0.0
    if not result:
        return 0.0
    if len(result) == len(expected_subjects):
        matches = 0
        for (res_subj, exp_subj) in zip(result, expected_subjects):
            if res_subj == exp_subj:
                matches += 1
        if matches == len(expected_subjects):
            return 1.0
    matches = 0
    result_idx = 0
    for exp_subj in expected_subjects:
        found = False
        for i in range(result_idx, len(result)):
            if result[i] == exp_subj:
                matches += 1
                result_idx = i + 1
                found = True
                break
        if not found:
            break
    return matches / len(expected_subjects)

def check_thunderbird_imap_settings__bd92e86076085c3fd0cc432ab151ca8f(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if Thunderbird IMAP settings are configured correctly.

    Args:
        result: Path to prefs.js file
        expected: Expected configuration with keys:
            - chunk_size: expected int value
            - min_chunk_size: expected int value
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    pref_pattern = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    found_chunk = False
    found_min_chunk = False
    chunk_matches = False
    min_chunk_matches = False
    with open(result, 'r') as f:
        for line in f:
            match = pref_pattern.match(line.strip())
            if match is None:
                continue
            key = match.group('key')
            value = json.loads(match.group('val'))
            if key == 'mail.imap.chunk_size':
                found_chunk = True
                chunk_matches = value == expected.get('chunk_size', 106496)
            elif key == 'mail.imap.min_chunk_size_threshold':
                found_min_chunk = True
                min_chunk_matches = value == expected.get('min_chunk_size', 159744)
    score = 0.0
    if found_chunk and chunk_matches:
        score += 0.5
    if found_min_chunk and min_chunk_matches:
        score += 0.5
    return score

def check_thunderbird_subject_date__2284321c(result, expected, **options):
    """Compare extracted subject and date data against expected values.

    Args:
        result: List of tuples (subject, date) from getter
        expected: Dict with 'expected_data' list of tuples
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_data = expected.get('expected_data', [])
    if len(result) != len(expected_data):
        return 0.0
    score = 0.0
    for (res_row, exp_row) in zip(result, expected_data):
        if res_row == tuple(exp_row):
            score += 1.0
    return score / len(expected_data) if expected_data else 0.0

def check_tb_filtered_emails__01abb20597fbc4723613687c7c60c0d3(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if the filtered email list matches expected criteria.

    Args:
        result: List of email dicts from getter
        expected: Dict with validation rules (e.g., 'keyword', 'min_count', 'expected_subjects')
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error('Invalid result format')
        return 0.0
    keyword = expected.get('keyword', '').lower()
    min_count = expected.get('min_count', 0)
    expected_subjects = expected.get('expected_subjects', [])
    score = 0.0
    if len(result) >= min_count:
        score += 0.3
    else:
        logger.debug(f'Count check failed: got {len(result)}, expected >= {min_count}')
    if keyword:
        keyword_match_count = sum((1 for email in result if keyword in email.get('subject', '').lower()))
        if keyword_match_count == len(result) and len(result) > 0:
            score += 0.4
        else:
            logger.debug(f"Keyword check: {keyword_match_count}/{len(result)} emails contain '{keyword}'")
    if expected_subjects:
        result_subjects = [email.get('subject', '').strip() for email in result]
        expected_subjects_lower = [s.lower() for s in expected_subjects]
        result_subjects_lower = [s.lower() for s in result_subjects]
        matches = sum((1 for exp_subj in expected_subjects_lower if exp_subj in result_subjects_lower))
        score += 0.3 * (matches / len(expected_subjects))
    return score

def check_thunderbird_three_cols__07170c7f(result, expected, **options):
    """Compare three-column data against expected values.

    Args:
        result: List of tuples (sender, subject, has_cc) from getter
        expected: Dict with 'expected_data' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_data = expected.get('expected_data', [])
    if len(result) != len(expected_data):
        return 0.0
    matches = 0
    for (res_row, exp_row) in zip(result, expected_data):
        if res_row == tuple(exp_row):
            matches += 1
    return matches / len(expected_data) if expected_data else 0.0

def check_thunderbird_txt_files__dd834e2c641c360491704c9db4c6f362(result, expected, **options):
    """
    Check if the directory listing contains expected .txt files matching patterns.

    Args:
        result: Path to file containing ls -R output
        expected: Expected patterns (dict with 'expect' list of regex patterns)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_patterns = [re.compile(ptt) for ptt in expected.get('expect', [])]
    if not expect_patterns:
        return 0.0
    expect_metrics = [False] * len(expect_patterns)
    with open(result, 'r') as f:
        content = f.read()
        for (i, pattern) in enumerate(expect_patterns):
            if pattern.search(content):
                expect_metrics[i] = True
    return float(all(expect_metrics))

def check_thunderbird_server_check__734b4908962171d93d8cfb6878a651d2(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if Thunderbird server settings are configured correctly.

    Args:
        result: Path to prefs.js file
        expected: Expected configuration with keys:
            - check_new_mail: expected boolean value
            - max_cached_connections: expected int value
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    pref_pattern = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    found_check = False
    found_connections = False
    check_matches = False
    connections_matches = False
    with open(result, 'r') as f:
        for line in f:
            match = pref_pattern.match(line.strip())
            if match is None:
                continue
            key = match.group('key')
            value = json.loads(match.group('val'))
            if key == 'mail.server.server1.check_new_mail':
                found_check = True
                check_matches = value == expected.get('check_new_mail', False)
            elif key == 'mail.server.server1.max_cached_connections':
                found_connections = True
                connections_matches = value == expected.get('max_cached_connections', 5)
    score = 0.0
    if found_check and check_matches:
        score += 0.5
    if found_connections and connections_matches:
        score += 0.5
    return score

def check_thunderbird_timeout__0d04b8a7(result: str, expected: dict) -> float:
    """
    Check if server timeout is set to expected value.

    Args:
        result: path to prefs.js file
        expected: dict with 'timeout' key (int: timeout in seconds)

    Returns:
        float: 1.0 if timeout matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_timeout = expected.get('timeout')
    if expected_timeout is None:
        logger.warning('No expected timeout value provided')
        return 0.0
    try:
        with open(result, 'r') as f:
            for line in f:
                if 'mail.server.server1.timeout' in line and 'user_pref' in line:
                    parts = line.split(', ')
                    if len(parts) >= 2:
                        value_part = parts[1].strip().rstrip(');')
                        actual_timeout = json.loads(value_part)
                        logger.debug(f'Found timeout: {actual_timeout}, expected: {expected_timeout}')
                        return 1.0 if actual_timeout == expected_timeout else 0.0
        logger.warning('mail.server.server1.timeout preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_email_deleted_and_folder_created__3c15b6d0(result: dict, expected: dict, **options):
    """Check if email was deleted and folder was created.

    Args:
        result: Dict from getter with email_count and folder_exists
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_email_count = expected.get('email_count', 2)
    expected_folder_exists = expected.get('folder_exists', True)
    score = 0.0
    if result.get('email_count') == expected_email_count:
        score += 0.5
        logger.info(f"Email count check passed: {result.get('email_count')}")
    else:
        logger.info(f"Email count check failed: got {result.get('email_count')}, expected {expected_email_count}")
    if result.get('folder_exists') == expected_folder_exists:
        score += 0.5
        logger.info('Folder existence check passed')
    else:
        logger.info(f"Folder existence check failed: got {result.get('folder_exists')}, expected {expected_folder_exists}")
    return score

def check_thunderbird_date_from_subject__a92c1922(result, expected, **options):
    """Compare extracted dates against expected values.

    Args:
        result: List of date strings from getter
        expected: Dict with 'expected_dates' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_dates = expected.get('expected_dates', [])
    if len(result) != len(expected_dates):
        return 0.0
    result_set = set(result)
    expected_set = set(expected_dates)
    if result_set == expected_set:
        return 1.0
    matches = len(result_set & expected_set)
    return matches / len(expected_dates) if expected_dates else 0.0

def check_thunderbird_check_new_mail__625e2b37(result: str, expected: dict) -> float:
    """
    Check if automatic mail checking is disabled for the account.

    Args:
        result: path to prefs.js file
        expected: dict with 'check_new_mail' key (bool: true/false)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_value = expected.get('check_new_mail')
    if expected_value is None:
        logger.warning('No expected check_new_mail value provided')
        return 0.0
    try:
        with open(result, 'r') as f:
            for line in f:
                if 'mail.server.server1.check_new_mail' in line and 'user_pref' in line:
                    parts = line.split(', ')
                    if len(parts) >= 2:
                        value_part = parts[1].strip().rstrip(');')
                        actual_value = json.loads(value_part)
                        logger.debug(f'Found check_new_mail: {actual_value}, expected: {expected_value}')
                        return 1.0 if actual_value == expected_value else 0.0
        logger.warning('mail.server.server1.check_new_mail preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_filter__2de049a0(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_email_domains__8d2fdf6f2182283d571d666af9c9e85f(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if the email domains match expected list.

    Args:
        result: List of email domains from the getter
        expected: Expected data with 'domains' key containing list of domain strings
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_domains = expected.get('domains', [])
    if not expected_domains:
        return 0.0
    result_normalized = [d.lower().strip() for d in result if d]
    expected_normalized = [d.lower().strip() for d in expected_domains]
    matches = 0
    for exp_domain in expected_normalized:
        if exp_domain in result_normalized:
            matches += 1
    score = matches / len(expected_normalized)
    return score

def check_thunderbird_workflow_and_amazon_tab__0f025d55a86976e1ae415a210c74f719(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Comprehensive verification of the Thunderbird-to-Chrome workflow.

    This metric verifies:
    1. Thunderbird Bills folder was accessed (from accessibility tree)
    2. Amazon.com was recently visited in Chrome (from history with timestamp)
    3. Exactly the expected tabs are open in Chrome (tab count and URLs)

    Args:
        result: Dictionary from get_thunderbird_and_chrome_state getter containing:
            {
                'thunderbird_state': {
                    'bills_folder_visible': bool,
                    'email_displayed': bool,
                    'folder_name': str or None
                },
                'chrome_history': [
                    {'url': str, 'title': str, 'last_visit_time': int},
                    ...
                ],
                'chrome_tabs': [
                    {'url': str, 'title': str},
                    ...
                ]
            }
        expected: Dictionary containing:
            {
                'exact_tab_urls': [list of expected tab URLs],
                'amazon_url': 'https://www.amazon.com/',
                'thunderbird_folder': 'Bills'
            }

    Returns:
        float: Score from 0.0 to 1.0
            - 1.0: All verifications passed (Thunderbird Bills accessed + recent Amazon visit + correct tabs)
            - 0.7: Tabs correct but Thunderbird verification failed (possible cheating)
            - 0.5: Partial verification (some checks passed)
            - 0.0: Failed verification
    """
    logger.info('=' * 80)
    logger.info('[THUNDERBIRD_WORKFLOW_CHECK] Starting comprehensive workflow verification')
    logger.info('=' * 80)
    if not result:
        logger.error('[THUNDERBIRD_WORKFLOW_CHECK] Result is None or empty')
        return 0.0
    thunderbird_state = result.get('thunderbird_state', {})
    chrome_history = result.get('chrome_history', [])
    chrome_tabs = result.get('chrome_tabs', [])
    expected_tab_urls = expected.get('exact_tab_urls', [])
    expected_amazon_url = expected.get('amazon_url', 'https://www.amazon.com/')
    expected_folder = expected.get('thunderbird_folder', 'Bills')
    logger.info(f'[THUNDERBIRD_WORKFLOW_CHECK] Expected tab URLs: {expected_tab_urls}')
    logger.info(f'[THUNDERBIRD_WORKFLOW_CHECK] Expected Amazon URL: {expected_amazon_url}')
    logger.info(f'[THUNDERBIRD_WORKFLOW_CHECK] Expected Thunderbird folder: {expected_folder}')
    checks = {'thunderbird_bills_folder': False, 'amazon_in_recent_history': False, 'correct_tab_count': False, 'correct_tab_urls': False}
    logger.info('-' * 80)
    logger.info('[CHECK 1] Verifying Thunderbird Bills folder access')
    logger.info(f'[CHECK 1] Thunderbird state: {thunderbird_state}')
    bills_folder_visible = thunderbird_state.get('bills_folder_visible', False)
    folder_name = thunderbird_state.get('folder_name', None)
    if bills_folder_visible:
        checks['thunderbird_bills_folder'] = True
        logger.info(f'[CHECK 1] ✓ PASSED - Bills folder is visible in Thunderbird')
        if folder_name == expected_folder:
            logger.info(f'[CHECK 1] ✓ BONUS - Bills folder is currently selected')
    else:
        logger.warning(f'[CHECK 1] ✗ FAILED - Bills folder not found in Thunderbird accessibility tree')
        logger.warning(f'[CHECK 1] This suggests the agent may not have accessed Thunderbird at all')
    logger.info('-' * 80)
    logger.info('[CHECK 2] Verifying Amazon.com in recent Chrome history')
    logger.info(f'[CHECK 2] Total history entries retrieved: {len(chrome_history)}')
    current_time = int(time.time())
    recency_threshold = 600
    amazon_found = False
    most_recent_amazon_time = 0
    most_recent_amazon_readable = 'Never'
    for entry in chrome_history:
        url = entry.get('url', '')
        visit_time = entry.get('last_visit_time', 0)
        readable_time = entry.get('visit_time_readable', 'Unknown')
        if 'amazon.com' in url.lower():
            amazon_found = True
            time_diff = current_time - visit_time
            logger.info(f'[CHECK 2] Found Amazon visit: {url}')
            logger.info(f'[CHECK 2] Visit time: {readable_time} (Unix: {visit_time})')
            logger.info(f'[CHECK 2] Time difference from now: {time_diff} seconds')
            if visit_time > most_recent_amazon_time:
                most_recent_amazon_time = visit_time
                most_recent_amazon_readable = readable_time
            if time_diff <= recency_threshold:
                checks['amazon_in_recent_history'] = True
                logger.info(f'[CHECK 2] ✓ PASSED - Amazon.com visited recently ({time_diff}s ago)')
                break
    if amazon_found and (not checks['amazon_in_recent_history']):
        time_diff = current_time - most_recent_amazon_time
        logger.warning(f'[CHECK 2] ✗ FAILED - Amazon.com found in history but visit is too old')
        logger.warning(f'[CHECK 2] Most recent Amazon visit: {most_recent_amazon_readable} ({time_diff}s ago)')
        logger.warning(f'[CHECK 2] Threshold: {recency_threshold}s')
    elif not amazon_found:
        logger.warning(f'[CHECK 2] ✗ FAILED - Amazon.com not found in recent Chrome history')
        logger.warning(f'[CHECK 2] This suggests the link may not have been clicked')
    logger.info('-' * 80)
    logger.info('[CHECK 3] Verifying Chrome tabs')
    logger.info(f'[CHECK 3] Expected {len(expected_tab_urls)} tabs with URLs: {expected_tab_urls}')
    if not chrome_tabs:
        logger.error('[CHECK 3] ✗ FAILED - No Chrome tabs found')
    else:
        actual_urls = [tab.get('url', '') for tab in chrome_tabs]
        logger.info(f'[CHECK 3] Found {len(actual_urls)} tabs with URLs: {actual_urls}')
        if len(actual_urls) == len(expected_tab_urls):
            checks['correct_tab_count'] = True
            logger.info(f'[CHECK 3] ✓ Tab count correct: {len(actual_urls)}')
        else:
            logger.warning(f'[CHECK 3] ✗ Tab count mismatch: expected {len(expected_tab_urls)}, got {len(actual_urls)}')
        all_urls_present = True
        for expected_url in expected_tab_urls:
            found = False
            for actual_url in actual_urls:
                if compare_urls(expected_url, actual_url):
                    found = True
                    break
            if not found:
                all_urls_present = False
                logger.warning(f'[CHECK 3] ✗ Expected URL not found: {expected_url}')
        no_unexpected_urls = True
        for actual_url in actual_urls:
            found = False
            for expected_url in expected_tab_urls:
                if compare_urls(actual_url, expected_url):
                    found = True
                    break
            if not found:
                no_unexpected_urls = False
                logger.warning(f'[CHECK 3] ✗ Unexpected URL found: {actual_url}')
        if all_urls_present and no_unexpected_urls:
            checks['correct_tab_urls'] = True
            logger.info(f'[CHECK 3] ✓ PASSED - All expected tabs present, no unexpected tabs')
        else:
            logger.warning(f'[CHECK 3] ✗ FAILED - Tab URL verification failed')
    logger.info('=' * 80)
    logger.info('[FINAL SCORING] Verification results:')
    logger.info(f"  - Thunderbird Bills folder accessed: {checks['thunderbird_bills_folder']}")
    logger.info(f"  - Amazon in recent history: {checks['amazon_in_recent_history']}")
    logger.info(f"  - Correct tab count: {checks['correct_tab_count']}")
    logger.info(f"  - Correct tab URLs: {checks['correct_tab_urls']}")
    if all(checks.values()):
        score = 1.0
        logger.info('[FINAL SCORING] ✓✓✓ ALL CHECKS PASSED - Task completed successfully')
        logger.info('[FINAL SCORING] Agent followed the complete Thunderbird workflow')
    elif checks['correct_tab_count'] and checks['correct_tab_urls']:
        if not checks['thunderbird_bills_folder']:
            score = 0.7
            logger.warning('[FINAL SCORING] ⚠ PARTIAL PASS (0.7) - Tabs correct but Thunderbird not accessed')
            logger.warning('[FINAL SCORING] Agent may have cheated by directly opening Amazon in Chrome')
        else:
            score = 0.8
            logger.warning('[FINAL SCORING] ⚠ PARTIAL PASS (0.8) - Thunderbird accessed, tabs correct, but recent history check failed')
            logger.warning('[FINAL SCORING] This might be a timing issue with history recording')
    elif checks['thunderbird_bills_folder']:
        score = 0.3
        logger.warning('[FINAL SCORING] ⚠ PARTIAL PASS (0.3) - Thunderbird accessed but tab verification failed')
    else:
        score = 0.0
        logger.error('[FINAL SCORING] ✗✗✗ VERIFICATION FAILED - Task not completed correctly')
    logger.info(f'[FINAL SCORING] Final score: {score}')
    logger.info('=' * 80)
    return score

def check_thunderbird_oldest_three__0c7bac1c(result, expected, **options):
    """Compare oldest 3 emails data against expected values with Thunderbird verification.

    Args:
        result: Dict with 'headers', 'data', and 'thunderbird_emails' keys from getter
                - headers: tuple of column headers
                - data: list of tuples (sender, subject) from Excel
                - thunderbird_emails: list of tuples (sender, subject, date) from Thunderbird daily folder
        expected: Dict with 'expected_data' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    headers = result.get('headers')
    data = result.get('data', [])
    thunderbird_emails = result.get('thunderbird_emails', [])
    expected_data = expected.get('expected_data', [])
    if headers is None or not isinstance(headers, tuple) or len(headers) != 2:
        return 0.0
    if headers[0] is None and headers[1] is None:
        return 0.0
    header_str = ' '.join([str(h).lower() if h else '' for h in headers])
    has_sender_like = any((keyword in header_str for keyword in ['sender', 'from', 'name']))
    has_subject_like = any((keyword in header_str for keyword in ['subject', 'title', 'topic']))
    header_score = 1.0 if has_sender_like and has_subject_like else 0.8
    if len(data) < len(expected_data):
        return 0.0
    if len(data) > len(expected_data):
        data = data[:len(expected_data)]
    temporal_verification_score = 1.0
    if thunderbird_emails and len(thunderbird_emails) >= 3:
        oldest_three_from_thunderbird = thunderbird_emails[:3]
        thunderbird_pairs = [(sender, subject) for (sender, subject, date) in oldest_three_from_thunderbird]

        def normalize_pair(pair):
            if isinstance(pair, (list, tuple)) and len(pair) >= 2:
                return (str(pair[0]).strip() if pair[0] else '', str(pair[1]).strip() if pair[1] else '')
            return ('', '')
        normalized_excel_data = [normalize_pair(row) for row in data]
        normalized_thunderbird_pairs = [normalize_pair(pair) for pair in thunderbird_pairs]
        matches_thunderbird = 0
        for excel_row in normalized_excel_data:
            if excel_row in normalized_thunderbird_pairs:
                matches_thunderbird += 1
        if matches_thunderbird < len(expected_data):
            temporal_verification_score = 0.0
        else:
            temporal_verification_score = 1.0
    else:
        temporal_verification_score = 0.7
    matches = 0
    for (res_row, exp_row) in zip(data, expected_data):
        if isinstance(res_row, tuple) and res_row == tuple(exp_row):
            matches += 1
        elif isinstance(res_row, list) and res_row == exp_row:
            matches += 1
    data_score = matches / len(expected_data) if expected_data else 0.0
    final_score = data_score * 0.5 + data_score * header_score * 0.1 + temporal_verification_score * 0.4
    return final_score

def check_email_count__fcf6b1e44de70adba5822f1fa1e262b7(result, expected, **options):
    """
    Check if the email count matches expected value or is within range.

    Args:
        result: Content of count file (string)
        expected: Expected count parameters (dict with 'min_count' or 'exact_count')
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    try:
        count = int(result.strip())
    except (ValueError, AttributeError):
        return 0.0
    if 'exact_count' in expected:
        return 1.0 if count == expected['exact_count'] else 0.0
    if 'min_count' in expected:
        return 1.0 if count >= expected['min_count'] else 0.0
    return 0.0

def check_contact_emails__0fa83097089d9008d43d2c029fbbd194(result, expected, **options):
    """Check if specific email addresses exist in the contact list.

    Args:
        result: List of actual email addresses from getter
        expected: Dict with 'required_emails' list from rules
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on how many required emails are found
    """
    required_emails = expected.get('required_emails', [])
    if not required_emails:
        return 0.0
    matches = sum((1 for email in required_emails if email in result))
    score = matches / len(required_emails)
    return score

def check_thunderbird_filter__9b7bc335(result: str, expected: Dict[str, List[Dict[str, Union[str, List[str]]]]]) -> float:
    """
    Check Thunderbird filter configuration for attachment-based copy filters.

    This function validates that a Thunderbird message filter correctly implements
    the requirement to copy emails with attachments to a specific folder.

    Validation checks:
    1. Filter is enabled (enabled="yes")
    2. Action type is correct (action="Copy to folder")
    3. Destination folder is specified (actionValue contains folder name)
    4. Attachment detection condition is present (flexible matching for various methods)

    The function accepts multiple valid attachment detection methods:
    - Size-based filtering (size,isGreaterThan)
    - Header inspection (Content-Disposition, Content-Type)
    - Attachment status fields (attachmentStatus, if available)
    - Multipart message detection

    Args:
        result (str): Path to Thunderbird filter definition file (msgFilterRules.dat)
        expected (Dict[str, List[Dict[str, Union[str, List[str]]]]]): Expected filter rules
            {
                "expect": [{
                    "enabled": "yes",
                    "action": "Copy to folder",
                    "actionValue": <folder_name_or_pattern>,
                    "condition": ["check_attachments"]  # Special marker for flexible matching
                }],
                "unexpect": [{...}]  # optional
            }

    Returns:
        float: 1.0 if all expected filters match and no unexpected filters found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        from desktop_env.evaluators.metrics.utils import _match_record
    except ImportError:
        from ..metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            if _match_filter_with_folder(r, flt):
                expect_metrics[i] = True
        for r in expected.get('unexpect', []):
            if _match_filter_with_folder(r, flt):
                unexpect_metric = False
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_prefs__249498a7cd34d43ebd7ef238b7ba6bee(result: str, expected: Dict[str, Dict[str, Dict[str, Any]]]) -> float:
    """
    Check Thunderbird preferences against expected rules.

    Args:
        result: Path to prefs.js file
        expected: Rules dict with 'expect' and optionally 'unexpect' keys

    Returns:
        1.0 if all expected preferences match and no unexpected preferences are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_rules = expected.get('expect', {})
    unexpect_rules = expected.get('unexpect', {})
    pref_pattern: Pattern[str] = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            match_: Match[str] = pref_pattern.match(line.strip())
            if match_ is None:
                continue
            key: str = match_.group('key')
            value = json.loads(match_.group('val'))
            if key in expect_rules:
                rule = expect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    expect_metrics[key] = value == ref
                elif method == 're':
                    expect_metrics[key] = bool(re.search(str(ref), str(value)))
                elif method == 'gt':
                    expect_metrics[key] = value > ref
                elif method == 'lt':
                    expect_metrics[key] = value < ref
                elif method == 'gte':
                    expect_metrics[key] = value >= ref
                elif method == 'lte':
                    expect_metrics[key] = value <= ref
            elif key in unexpect_rules:
                rule = unexpect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    unexpect_metric = unexpect_metric and value != ref
                elif method == 're':
                    unexpect_metric = unexpect_metric and (not bool(re.search(str(ref), str(value))))
    return float(all(expect_metrics.values()) and unexpect_metric)

def check_thunderbird_drafts_picker__17405a90(result: str, expected: dict) -> float:
    """
    Check if drafts folder picker mode matches expected value.

    Args:
        result: path to prefs.js file
        expected: dict with 'drafts_folder_picker_mode' key

    Returns:
        float: 1.0 if mode matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_mode = expected.get('drafts_folder_picker_mode')
    if expected_mode is None:
        logger.warning('No expected drafts_folder_picker_mode provided')
        return 0.0
    try:
        with open(result, 'r') as f:
            for line in f:
                if 'mail.identity.id1.drafts_folder_picker_mode' in line and 'user_pref' in line:
                    parts = line.split('"')
                    if len(parts) >= 4:
                        actual_mode = parts[3]
                        logger.debug(f'Found drafts_folder_picker_mode: {actual_mode}, expected: {expected_mode}')
                        return 1.0 if actual_mode == expected_mode else 0.0
        logger.warning('mail.identity.id1.drafts_folder_picker_mode preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_filter__ce00a590(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_record_with_regex(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record_with_regex(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_filter__07f67672(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_folder_exists_and_email_starred__7e32e736(result: dict, expected: dict, **options):
    """Check if folder exists and email is starred.

    Args:
        result: Dict from getter with folder_exists and has_starred_email
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_folder_exists = expected.get('folder_exists', True)
    expected_has_starred = expected.get('has_starred_email', True)
    score = 0.0
    if result.get('folder_exists') == expected_folder_exists:
        score += 0.5
        logger.info('Folder existence check passed')
    else:
        logger.info(f"Folder existence check failed: got {result.get('folder_exists')}, expected {expected_folder_exists}")
    if result.get('has_starred_email') == expected_has_starred:
        score += 0.5
        logger.info('Starred email check passed')
    else:
        logger.info(f"Starred email check failed: got {result.get('has_starred_email')}, expected {expected_has_starred}")
    return score

def check_thunderbird_cc_emails__f218f3c2(result, expected, **options):
    """Compare CC email data against expected values.

    Verifies that:
    1. Excel file was created with correct data
    2. Data in Excel matches the CC-filtered emails from Thunderbird
    3. The emails match the expected values

    Args:
        result: Dict from getter with:
            - 'excel_data': List of tuples (sender, subject) from Excel
            - 'thunderbird_cc_emails': List of tuples from Thunderbird
            - 'valid': bool indicating if Excel matches Thunderbird
        expected: Dict with 'expected_data' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.warning(f'Invalid types: result={type(result)}, expected={type(expected)}')
        return 0.0
    expected_data = expected.get('expected_data', [])
    if not expected_data:
        logger.warning('No expected data provided')
        return 0.0
    excel_data = result.get('excel_data', [])
    thunderbird_data = result.get('thunderbird_cc_emails', [])
    is_valid = result.get('valid', False)
    if not is_valid:
        logger.warning('Excel data does not match Thunderbird emails - user may have manually typed values')
        return 0.0
    expected_tuples = [tuple(row) for row in expected_data]
    if len(excel_data) != len(expected_tuples):
        logger.warning(f'Length mismatch: got {len(excel_data)}, expected {len(expected_tuples)}')
        return 0.0
    matches = 0
    for (res_row, exp_row) in zip(excel_data, expected_tuples):
        if res_row == exp_row:
            matches += 1
        else:
            logger.debug(f'Row mismatch: got {res_row}, expected {exp_row}')
    score = matches / len(expected_tuples) if expected_tuples else 0.0
    logger.info(f'Score: {score} ({matches}/{len(expected_tuples)} rows matched)')
    return score

def check_thunderbird_sender_count__44f17dcf(result, expected, **options):
    """Compare sender count data against expected values with dual verification.

    This metric performs TWO levels of verification:
    1. Verifies that Thunderbird mailbox data matches expected sender counts
    2. Verifies that Excel report matches the Thunderbird source data

    This prevents users from manually creating an Excel file with expected values
    without actually analyzing Thunderbird emails.

    Args:
        result: Dict from getter with keys:
            - 'excel_sender_counts': sender counts from Excel file
            - 'thunderbird_sender_counts': sender counts from Thunderbird mailbox
            - 'excel_valid': bool indicating Excel was parsed successfully
            - 'thunderbird_valid': bool indicating Thunderbird was parsed successfully
        expected: Dict with 'expected_counts' mapping sender names to counts
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_counts = expected.get('expected_counts', {})
    if not expected_counts:
        return 0.0
    excel_counts = result.get('excel_sender_counts', {})
    thunderbird_counts = result.get('thunderbird_sender_counts', {})
    excel_valid = result.get('excel_valid', False)
    thunderbird_valid = result.get('thunderbird_valid', False)
    if not excel_valid or not thunderbird_valid:
        return 0.0
    total_points = 0.0
    max_points = 3.0
    thunderbird_matches = 0
    for (sender, count) in expected_counts.items():
        if sender in thunderbird_counts and thunderbird_counts[sender] == count:
            thunderbird_matches += 1
    if thunderbird_matches == len(expected_counts):
        total_points += 1.0
    else:
        total_points += thunderbird_matches / len(expected_counts)
    excel_matches = 0
    for (sender, count) in expected_counts.items():
        if sender in excel_counts and excel_counts[sender] == count:
            excel_matches += 1
    if excel_matches == len(expected_counts):
        total_points += 1.0
    else:
        total_points += excel_matches / len(expected_counts)
    consistency_score = calculate_consistency_score(excel_counts, thunderbird_counts, expected_counts)
    total_points += consistency_score
    return total_points / max_points

def check_all_emails_starred__d14b39d6(result, expected, **options):
    """Check if all emails in a specific folder are starred.

    Args:
        result: Path to the SQLite database file
        expected: Dict with 'folder_id' key specifying the folder
        **options: Additional options

    Returns:
        float: 1.0 if all emails are starred, 0.0 otherwise
    """
    import sqlite3
    folder_id = expected.get('folder_id')
    if folder_id is None:
        return 0.0
    connection = sqlite3.connect(result)
    cursor = connection.cursor()
    cursor.execute('SELECT COUNT(*) FROM messages WHERE folderID = ?', (folder_id,))
    total_count = cursor.fetchone()[0]
    if total_count == 0:
        connection.close()
        return 0.0
    cursor.execute('\n        SELECT COUNT(*) FROM messageAttributes\n        WHERE attributeID = 58 AND value = 1\n        AND messageID IN (SELECT id FROM messages WHERE folderID = ?)\n    ', (folder_id,))
    starred_count = cursor.fetchone()[0]
    connection.close()
    return 1.0 if starred_count == total_count else 0.0

def check_thunderbird_filter_created__8a2db6ad(result: str, expected: dict) -> float:
    """
    Check if a message filter with expected properties was created.

    Args:
        result: path to msgFilterRules.dat file
        expected: dict with 'filter_name', 'sender', and 'target_folder' keys

    Returns:
        float: Score based on how many properties match (0.0 to 1.0)
    """
    if result is None:
        return 0.0
    expected_name = expected.get('filter_name')
    expected_sender = expected.get('sender')
    expected_folder = expected.get('target_folder')
    if not expected_name:
        logger.warning('No expected filter_name provided')
        return 0.0
    try:
        with open(result, 'r') as f:
            lines = f.readlines()
        filters = []
        current_filter = None
        for line in lines:
            line = line.strip()
            if not line:
                continue
            if line.startswith('name='):
                if current_filter is not None:
                    filters.append(current_filter)
                current_filter = {}
                current_filter['name'] = _value_processor(line[6:-1].strip('"'))
            elif current_filter is not None:
                if line.startswith('enabled='):
                    current_filter['enabled'] = _value_processor(line[9:-1].strip('"'))
                elif line.startswith('type='):
                    current_filter['type'] = _value_processor(line[6:-1].strip('"'))
                elif line.startswith('action='):
                    current_filter['action'] = _value_processor(line[8:-1].strip('"'))
                elif line.startswith('actionValue='):
                    current_filter['actionValue'] = _value_processor(line[13:-1].strip('"'))
                elif line.startswith('condition='):
                    current_filter['condition'] = _value_processor(line[11:-1].strip('"'))
                    filters.append(current_filter)
                    current_filter = None
        if current_filter is not None:
            filters.append(current_filter)
        logger.debug(f'Parsed {len(filters)} filters from msgFilterRules.dat')
        matching_filter = None
        for flt in filters:
            if flt.get('name') == expected_name:
                matching_filter = flt
                break
        if matching_filter is None:
            logger.warning(f"Filter '{expected_name}' not found in msgFilterRules.dat")
            return 0.0
        logger.debug(f'Found filter: {matching_filter}')
        score = 0.0
        score += 0.33
        logger.debug(f'Filter name matched: {expected_name}')
        if expected_sender:
            condition = matching_filter.get('condition', '')
            if expected_sender in condition:
                score += 0.33
                logger.debug(f"Sender '{expected_sender}' found in condition: {condition}")
            else:
                logger.warning(f"Sender '{expected_sender}' not found in condition: {condition}")
        else:
            score += 0.33
        if expected_folder:
            action = matching_filter.get('action', '')
            action_value = matching_filter.get('actionValue', '')
            if 'Move to folder' in action:
                if expected_folder in action_value:
                    score += 0.34
                    logger.debug(f"Target folder '{expected_folder}' found in actionValue: {action_value}")
                else:
                    logger.warning(f"Target folder '{expected_folder}' not found in actionValue: {action_value}")
            else:
                logger.warning(f"Action is not 'Move to folder': {action}")
        else:
            score += 0.34
        logger.debug(f'Final score: {score}')
        return score
    except Exception as e:
        logger.error(f'Error reading msgFilterRules.dat: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_thunderbird_amazon_deleted__4bba48e3(result_state: List[Dict[str, str]], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if Amazon emails were deleted from Thunderbird folder.

    This function verifies that:
    1. All emails containing 'Amazon' in the subject have been deleted
    2. At least some non-Amazon emails remain in the folder

    Args:
        result_state: List of email dicts from getter, each with 'subject' and 'raw' keys
        expected_state: Dict with rules:
            - no_amazon (bool): If True, verifies no emails with 'Amazon' in subject remain
            - min_remaining (int): Minimum number of emails that should remain after deletion
        **options: Additional options

    Returns:
        float: 1.0 if verification passes (Amazon emails deleted, min emails remaining), 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    no_amazon = expected_state.get('no_amazon', True) if isinstance(expected_state, dict) else True
    min_remaining = expected_state.get('min_remaining', 1) if isinstance(expected_state, dict) else 1
    logger.info(f'Checking emails: no_amazon={no_amazon}, min_remaining={min_remaining}')
    logger.info(f'Found {len(result_state)} emails in folder')
    if len(result_state) < min_remaining:
        logger.warning(f'Not enough emails remaining: {len(result_state)} < {min_remaining}')
        return 0.0
    if no_amazon:
        amazon_count = 0
        for email in result_state:
            subject = email.get('subject', '')
            if 'amazon' in subject.lower():
                amazon_count += 1
                logger.warning(f'Found Amazon email that should have been deleted: {subject}')
        if amazon_count > 0:
            logger.warning(f'Found {amazon_count} Amazon email(s) that should have been deleted')
            return 0.0
    logger.info(f'Verification passed: {len(result_state)} emails remain, no Amazon emails found')
    return 1.0

def check_thunderbird_folder_and_empty_bills__a91026e2(result: dict, expected: dict, **options):
    """Check if folder exists and Bills folder is empty.

    Args:
        result: Dict from getter with folder_exists and bills_email_count
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_folder_exists = expected.get('folder_exists', True)
    expected_email_count = expected.get('bills_email_count', 0)
    score = 0.0
    if result.get('folder_exists') == expected_folder_exists:
        score += 0.5
        logger.info('Folder existence check passed')
    else:
        logger.info(f"Folder existence check failed: got {result.get('folder_exists')}, expected {expected_folder_exists}")
    if result.get('bills_email_count') == expected_email_count:
        score += 0.5
        logger.info(f"Bills email count check passed: {result.get('bills_email_count')}")
    else:
        logger.info(f"Bills email count check failed: got {result.get('bills_email_count')}, expected {expected_email_count}")
    return score

def check_first_email__105b8e17(result, expected, **options):
    """
    Check if first contact email matches expected.

    Args:
        result: str email from getter
        expected: str expected email

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    return 1.0 if result == expected else 0.0

def check_thunderbird_compact_folders__e06bf170da6e0042eda0ceb7f1f9f833(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if Thunderbird is configured to automatically prompt for folder compaction.

    Args:
        result: Path to prefs.js file
        expected: Expected configuration with keys:
            - prompt_purge_enabled: boolean - false means prompting IS enabled
                                    (double negative: false = do NOT skip prompt)
            - purge_threshold_mb: int - size threshold in MB for prompting
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0

    Note:
        Thunderbird uses 'threshhold' spelling (typo) in actual preference names.
        mail.prompt_purge_threshhold=false means prompting is enabled.
        mail.purge_threshhold_mb sets the size threshold.
    """
    if result is None:
        return 0.0
    pref_pattern = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    found_compact = False
    found_threshold = False
    compact_matches = False
    threshold_matches = False
    with open(result, 'r') as f:
        for line in f:
            match = pref_pattern.match(line.strip())
            if match is None:
                continue
            key = match.group('key')
            value = json.loads(match.group('val'))
            if key == 'mail.prompt_purge_threshhold':
                found_compact = True
                compact_matches = value == expected.get('prompt_purge_enabled', False)
            elif key == 'mail.purge_threshhold_mb':
                found_threshold = True
                threshold_matches = value == expected.get('purge_threshold_mb', 20)
    score = 0.0
    if found_compact and compact_matches:
        score += 0.5
    if found_threshold and threshold_matches:
        score += 0.5
    return score

def check_thunderbird_subject__c9ce3f52(result, expected, **options):
    """Check if the email subject contains the expected text.

    Args:
        result: The actual subject text from getter
        expected: Dict with 'text' key containing expected substring
        **options: Additional options (case_sensitive, exact_match)

    Returns:
        float: 1.0 if match found, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_text = expected.get('text', '')
    case_sensitive = options.get('case_sensitive', False)
    exact_match = options.get('exact_match', False)
    if not case_sensitive:
        result = result.lower()
        expected_text = expected_text.lower()
    if exact_match:
        return 1.0 if result == expected_text else 0.0
    else:
        return 1.0 if expected_text in result else 0.0

def check_tb_digest__b6ecd864d1dd0f24d97ae590f2c596e8(result, expected, **options):
    """Compare extracted digest emails against expected values.

    Args:
        result: List of email dicts from getter
        expected: Dict with 'emails' key containing expected email list

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_emails = expected.get('emails', [])
    if len(result) != len(expected_emails):
        return 0.0
    result_dates = []
    for email in result:
        date_str = email.get('date', '')
        if date_str:
            try:
                result_dates.append(datetime.strptime(date_str, '%Y-%m-%d'))
            except:
                result_dates.append(None)
        else:
            result_dates.append(None)
    for i in range(len(result_dates) - 1):
        if result_dates[i] is None or result_dates[i + 1] is None:
            continue
        if result_dates[i] > result_dates[i + 1]:
            return 0.0
    score = 0.0
    total_fields = len(expected_emails) * 3
    for exp_email in expected_emails:
        best_match_score = 0
        for res_email in result:
            matches = 0
            if res_email.get('sender_name') == exp_email.get('sender_name'):
                matches += 1
            if res_email.get('subject') == exp_email.get('subject'):
                matches += 1
            if res_email.get('date') == exp_email.get('date'):
                matches += 1
            if matches > best_match_score:
                best_match_score = matches
        score += best_match_score
    return score / total_fields

def check_tb_summary__f5f661c56b53d0549f4c6bd9201ba899(result, expected, **options):
    """Compare email count summary against expected values and validate all requirements.

    Args:
        result: Dict with keys:
            - 'summary': Dict mapping sender emails to counts
            - 'headers_valid': bool (True if headers are correct)
            - 'sorted_valid': bool (True if sorted by count descending)
            - 'location_valid': bool (True if file is on Desktop)
            - 'filename_valid': bool (True if filename is 'email_summary.xlsx')
        expected: Dict with 'summary' key containing expected counts

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    if not result.get('headers_valid', False):
        return 0.0
    if not result.get('sorted_valid', False):
        return 0.0
    if not result.get('location_valid', False):
        return 0.0
    if not result.get('filename_valid', False):
        return 0.0
    actual_summary = result.get('summary', {})
    expected_summary = expected.get('summary', {})
    if len(actual_summary) != len(expected_summary):
        return 0.0
    score = 0.0
    for (sender, count) in expected_summary.items():
        if actual_summary.get(sender) == count:
            score += 1.0 / len(expected_summary)
    return score

def check_email_column__3631675f(result, expected, **options):
    """Check if email column values match expected list.

    Args:
        result: List of actual email values
        expected: Dict with 'emails' list to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_emails = expected.get('emails', [])
    if not isinstance(result, list):
        return 0.0
    score = 0.0
    total = len(expected_emails)
    for (i, (actual, exp)) in enumerate(zip(result, expected_emails)):
        if actual == exp:
            score += 1.0 / total
        elif exp is None and actual is None:
            score += 1.0 / total
        elif isinstance(actual, str) and isinstance(exp, str):
            if actual.strip().lower() == exp.strip().lower():
                score += 1.0 / total
    return score

def check_email_subjects_list__df60c313c1f3b2e712e9756c47386ddc(result, expected, **options):
    """
    Check if the email subjects file contains all expected subject lines.
    Allows additional subjects beyond the expected ones.

    Args:
        result: Content of the subjects file (str)
        expected: Expected configuration (dict with 'expect_subjects' list and 'min_count')
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_subjects = expected.get('expect_subjects', [])
    min_count = expected.get('min_count', len(expect_subjects))
    if not expect_subjects:
        return 0.0
    result_lines = [line.strip() for line in result.split('\n') if line.strip()]
    if len(result_lines) < min_count:
        return 0.0
    found_subjects = [False] * len(expect_subjects)
    for line in result_lines:
        for (i, expected_subject) in enumerate(expect_subjects):
            if line == expected_subject:
                found_subjects[i] = True
    if not all(found_subjects):
        return 0.0
    return 1.0

def check_all_emails_read__20eb9ee4(result, expected, **options):
    """Check if all emails in a specific folder are marked as read.

    Args:
        result: Path to the SQLite database file
        expected: Dict with 'folder_id' key specifying the folder
        **options: Additional options

    Returns:
        float: 1.0 if all emails are read, 0.0 otherwise
    """
    import sqlite3
    folder_id = expected.get('folder_id')
    if folder_id is None:
        return 0.0
    connection = sqlite3.connect(result)
    cursor = connection.cursor()
    cursor.execute('SELECT COUNT(*) FROM messages WHERE folderID = ?', (folder_id,))
    total_count = cursor.fetchone()[0]
    if total_count == 0:
        connection.close()
        return 0.0
    cursor.execute('\n        SELECT COUNT(*) FROM messageAttributes\n        WHERE attributeID = 59 AND value = 1\n        AND messageID IN (SELECT id FROM messages WHERE folderID = ?)\n    ', (folder_id,))
    read_count = cursor.fetchone()[0]
    connection.close()
    return 1.0 if read_count == total_count else 0.0

def check_tb_sender_counts__6a3e0f1dfc20fc63f235dcd50f4dfaaf(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """
    Check if the sender counts in the spreadsheet match expected values.

    Args:
        result: Dict mapping sender names to counts from getter
        expected: Dict with 'sender_counts' key containing expected counts
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Invalid result format')
        return 0.0
    expected_counts = expected.get('sender_counts', {})
    if not expected_counts:
        logger.error('No expected counts provided')
        return 0.0
    correct_count = 0
    total_count = len(expected_counts)
    for (sender, expected_count) in expected_counts.items():
        if sender in result and result[sender] == expected_count:
            correct_count += 1
        else:
            logger.debug(f"Sender '{sender}': expected {expected_count}, got {result.get(sender, 0)}")
    return correct_count / total_count if total_count > 0 else 0.0

def check_thunderbird_trash_folder__550830d3(result: str, expected: dict) -> float:
    """
    Check if trash folder name matches expected value for the Outlook account.

    Args:
        result: path to prefs.js file
        expected: dict with 'trash_folder_name' key

    Returns:
        float: 1.0 if folder name matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_name = expected.get('trash_folder_name')
    if not expected_name:
        logger.warning('No expected trash_folder_name provided')
        return 0.0
    try:
        server_id = None
        server_to_username = {}
        identity_to_email = {}
        identity_to_server = {}
        with open(result, 'r') as f:
            for line in f:
                if 'mail.server.server' in line and '.userName' in line and ('user_pref' in line):
                    parts = line.split('"')
                    if len(parts) >= 4:
                        pref_key = parts[1]
                        username = parts[3]
                        if 'mail.server.server' in pref_key and '.userName' in pref_key:
                            server_part = pref_key.split('.')[2]
                            server_to_username[server_part] = username
                if 'mail.identity.id' in line and '.useremail' in line and ('user_pref' in line):
                    parts = line.split('"')
                    if len(parts) >= 4:
                        pref_key = parts[1]
                        email = parts[3]
                        if 'mail.identity.id' in pref_key and '.useremail' in pref_key:
                            identity_part = pref_key.split('.')[2]
                            identity_to_email[identity_part] = email
                if 'mail.identity.id' in line and '.smtpServer' in line and ('user_pref' in line):
                    parts = line.split('"')
                    if len(parts) >= 4:
                        pref_key = parts[1]
                        smtp_server = parts[3]
                        if 'mail.identity.id' in pref_key and '.smtpServer' in pref_key:
                            identity_part = pref_key.split('.')[2]
                            identity_to_server[identity_part] = smtp_server
        target_email = 'anonym-x2024@outlook.com'
        for (server, username) in server_to_username.items():
            if username == target_email:
                server_id = server
                logger.debug(f'Found server {server_id} for email {target_email} via userName')
                break
        if server_id is None:
            for (identity, email) in identity_to_email.items():
                if email == target_email:
                    logger.debug(f'Found identity {identity} for email {target_email}')
                    if len(server_to_username) == 1:
                        server_id = list(server_to_username.keys())[0]
                        logger.debug(f'Using single server {server_id}')
                    break
        if server_id is None:
            logger.warning(f'Could not find server ID for email {target_email}')
            return 0.0
        trash_pref_key = f'mail.server.{server_id}.trash_folder_name'
        logger.debug(f'Looking for preference: {trash_pref_key}')
        with open(result, 'r') as f:
            for line in f:
                if trash_pref_key in line and 'user_pref' in line:
                    parts = line.split('"')
                    if len(parts) >= 4:
                        actual_name = parts[3]
                        logger.debug(f'Found trash_folder_name: {actual_name}, expected: {expected_name}')
                        return 1.0 if actual_name == expected_name else 0.0
        logger.warning(f'{trash_pref_key} preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_emails__f5c13cdd(result: str, rules: List[Dict[str, Any]], osname: str='ubuntu') -> float:
    """
    Check Thunderbird email recipients with partial credit support.

    This function checks how many of the required email addresses are present in the
    To field of a Thunderbird compose window. Unlike check_accessibility_tree which
    returns 0.0 if any rule fails, this function implements true partial credit by
    counting how many rules pass and returning the ratio.

    Args:
        result (str): XML of GNOME Accessibility Tree
        rules (List[Dict[str, Any]]): list of dict with selectors or xpath for each email to check
        osname (str): "ubuntu" | "windows" | "macos". "ubuntu" by default.

    Returns:
        float: Partial credit score between 0.0 and 1.0 (count of passed rules / total rules)
    """
    a11y_ns_map = _accessibility_ns_map[osname]
    try:
        at: _Element = lxml.etree.fromstring(result)
    except Exception as e:
        logger.error(f'Failed to parse accessibility tree: {e}')
        return 0.0
    if not rules:
        logger.warning('No rules provided')
        return 0.0
    passed_count = 0
    total_count = len(rules)
    for (i, r) in enumerate(rules):
        try:
            if 'xpath' in r:
                elements: List[_Element] = at.xpath(r['xpath'], namespaces=a11y_ns_map)
            elif 'selectors' in r:
                selector = CSSSelector(', '.join(r['selectors']), namespaces=a11y_ns_map)
                elements: List[_Element] = selector(at)
            else:
                logger.warning(f'Rule {i} missing xpath and selectors')
                continue
            if len(elements) > 0:
                passed_count += 1
                logger.info(f'Rule {i} passed: found {len(elements)} matching elements')
            else:
                logger.info(f"Rule {i} failed: no matching elements for {r.get('xpath', r.get('selectors', 'unknown'))}")
        except Exception as e:
            logger.error(f'Error processing rule {i}: {e}')
            continue
    score = passed_count / total_count if total_count > 0 else 0.0
    logger.info(f'Partial credit score: {passed_count}/{total_count} = {score:.2f}')
    return float(score)

def check_thunderbird_filter__793e1ffa(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_tb_first3_emails__843811624c18ed6ed65e7d4d877f3122(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if the first 3 emails match expected values.

    Args:
        result: List of email dicts from getter
        expected: Dict with 'expected_emails' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error('Invalid result format')
        return 0.0
    expected_emails = expected.get('expected_emails', [])
    if len(result) != 3:
        logger.error(f'Expected 3 emails, got {len(result)}')
        return 0.0
    score = 0.0
    for (i, (exp_email, res_email)) in enumerate(zip(expected_emails, result)):
        email_score = 0.0
        if exp_email.get('sender_name', '').strip() == res_email.get('sender_name', '').strip():
            email_score += 0.33
        if exp_email.get('subject', '').strip() == res_email.get('subject', '').strip():
            email_score += 0.34
        exp_date = exp_email.get('date', '').strip()
        res_date = res_email.get('date', '').strip()
        if exp_date in res_date or res_date in exp_date or exp_date == res_date:
            email_score += 0.33
        score += email_score / 3.0
    return score

def check_thunderbird_all_addresses__d8fe40b3(result, expected, **options):
    """Compare sender and CC address data against expected values.

    Args:
        result: List of tuples (sender_address, cc_address) from getter
        expected: Dict with 'expected_data' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_data = expected.get('expected_data', [])
    if len(result) != len(expected_data):
        return 0.0
    matches = 0
    for (res_row, exp_row) in zip(result, expected_data):
        exp_tuple = (exp_row[0], exp_row[1] if len(exp_row) > 1 else None)
        if res_row == exp_tuple:
            matches += 1
    return matches / len(expected_data) if expected_data else 0.0

def check_thunderbird_pdfs_uploaded_to_gdrive__69a89619e7273cf1d986af62cad42166(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF attachments from Thunderbird email were uploaded to Google Drive.

    This metric validates:
    1. The email exists in Thunderbird
    2. The email has PDF attachments
    3. All PDF attachments are uploaded to Google Drive root folder

    Args:
        result: Dict with keys:
            - email_data: Dict from Thunderbird getter with keys:
                - pdf_attachments: List of PDF filenames from email
                - email_found: bool
                - attachment_count: int
            - gdrive_data: Dict from Google Drive getter with keys:
                - file_count: int
                - file_names: List of PDF filenames on Google Drive
        expected: Dict with validation rules (not used, validation is implicit)

    Returns:
        float: 1.0 if all PDFs uploaded correctly, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dict')
        return 0.0
    email_data = result.get('email_data', {})
    gdrive_data = result.get('gdrive_data', {})
    if not email_data.get('email_found', False):
        logger.error("Email 'Paper Recommendation' not found in Thunderbird")
        return 0.0
    email_pdfs = email_data.get('pdf_attachments', [])
    attachment_count = email_data.get('attachment_count', 0)
    if attachment_count == 0:
        logger.error('No PDF attachments found in the email')
        return 0.0
    gdrive_pdfs = gdrive_data.get('file_names', [])
    gdrive_count = gdrive_data.get('file_count', 0)
    if gdrive_count == 0:
        logger.error('No PDF files found in Google Drive root folder')
        return 0.0
    email_pdfs_normalized = {pdf.strip().lower() for pdf in email_pdfs}
    gdrive_pdfs_normalized = {pdf.strip().lower() for pdf in gdrive_pdfs}
    missing_pdfs = email_pdfs_normalized - gdrive_pdfs_normalized
    if missing_pdfs:
        logger.error(f'Missing PDFs in Google Drive: {missing_pdfs}')
        logger.info(f'Email PDFs: {email_pdfs}')
        logger.info(f'Google Drive PDFs: {gdrive_pdfs}')
        return 0.0
    logger.info(f'Successfully verified: {len(email_pdfs)} PDFs from email uploaded to Google Drive')
    return 1.0

def check_thunderbird_filter__e2c8667a(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_filter(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_filter(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_prefs__5b041e31e4566087bb3bbe0b347cd7ff(result: str, expected: Dict[str, Dict[str, Dict[str, Any]]]) -> float:
    """
    Check Thunderbird preferences against expected rules.

    Args:
        result: Path to prefs.js file
        expected: Rules dict with 'expect' and optionally 'unexpect' keys

    Returns:
        1.0 if all expected preferences match and no unexpected preferences are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_rules = expected.get('expect', {})
    unexpect_rules = expected.get('unexpect', {})
    pref_pattern: Pattern[str] = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            match_: Match[str] = pref_pattern.match(line.strip())
            if match_ is None:
                continue
            key: str = match_.group('key')
            value = json.loads(match_.group('val'))
            if key in expect_rules:
                rule = expect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    expect_metrics[key] = value == ref
                elif method == 're':
                    expect_metrics[key] = bool(re.search(str(ref), str(value)))
                elif method == 'gt':
                    expect_metrics[key] = value > ref
                elif method == 'lt':
                    expect_metrics[key] = value < ref
                elif method == 'gte':
                    expect_metrics[key] = value >= ref
                elif method == 'lte':
                    expect_metrics[key] = value <= ref
            elif key in unexpect_rules:
                rule = unexpect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    unexpect_metric = unexpect_metric and value != ref
                elif method == 're':
                    unexpect_metric = unexpect_metric and (not bool(re.search(str(ref), str(value))))
    return float(all(expect_metrics.values()) and unexpect_metric)

def check_thunderbird_reply_settings__616461c5ba008feaf990bbd287063304(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if Thunderbird reply settings are configured correctly.

    Args:
        result: Path to prefs.js file
        expected: Expected configuration with keys:
            - reply_on_top: expected int value (0=bottom, 1=top, 2=select)
            - do_bcc: expected boolean value
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    pref_pattern = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    found_reply = False
    found_bcc = False
    reply_matches = False
    bcc_matches = False
    with open(result, 'r') as f:
        for line in f:
            match = pref_pattern.match(line.strip())
            if match is None:
                continue
            key = match.group('key')
            value = json.loads(match.group('val'))
            if key == 'mail.identity.id1.reply_on_top':
                found_reply = True
                reply_matches = value == expected.get('reply_on_top', 1)
            elif key == 'mail.identity.id1.doBcc':
                found_bcc = True
                bcc_matches = value == expected.get('do_bcc', False)
    score = 0.0
    if found_reply and reply_matches:
        score += 0.5
    if found_bcc and bcc_matches:
        score += 0.5
    return score

def check_tb_weekly__264d715cc53b32b5aee947baaf88b841(result, expected, **options):
    """Compare extracted weekly emails against expected values.

    Args:
        result: Dict with 'emails' list, 'headers' list, and 'valid_headers' boolean
        expected: Dict with 'emails' list, 'date_range' dict, and 'required_headers' list

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    result_emails = result.get('emails', [])
    result_headers = result.get('headers', [])
    valid_headers = result.get('valid_headers', False)
    expected_emails = expected.get('emails', [])
    expected_headers = expected.get('required_headers', [])
    date_range = expected.get('date_range', {})
    score_components = []
    if valid_headers:
        score_components.append(0.2)
    else:
        score_components.append(0.0)
    if len(result_emails) == len(expected_emails):
        score_components.append(0.1)
    else:
        score_components.append(0.0)
    date_range_score = 0.0
    if date_range.get('start') and date_range.get('end'):
        start_date = datetime.strptime(date_range['start'], '%Y-%m-%d')
        end_date = datetime.strptime(date_range['end'], '%Y-%m-%d')
        all_dates_valid = True
        for email in result_emails:
            email_date_str = email.get('date')
            if email_date_str:
                try:
                    email_date = datetime.strptime(email_date_str, '%Y-%m-%d')
                    if not start_date <= email_date <= end_date:
                        all_dates_valid = False
                        break
                except ValueError:
                    all_dates_valid = False
                    break
            else:
                all_dates_valid = False
                break
        if all_dates_valid and len(result_emails) > 0:
            date_range_score = 0.2
    score_components.append(date_range_score)
    content_score = 0.0
    if len(expected_emails) > 0:
        expected_tuples = set()
        for exp_email in expected_emails:
            expected_tuples.add((exp_email.get('sender_name'), exp_email.get('sender_email'), exp_email.get('subject'), exp_email.get('date')))
        result_tuples = set()
        for res_email in result_emails:
            result_tuples.add((res_email.get('sender_name'), res_email.get('sender_email'), res_email.get('subject'), res_email.get('date')))
        matched = expected_tuples & result_tuples
        content_score = 0.5 * (len(matched) / len(expected_emails))
    score_components.append(content_score)
    return sum(score_components)

def check_name_email_table__6e2aeb6fcb42c45735613c9af17f7e4d(result: List[List[str]], expected: Dict[str, Any], **options) -> float:
    """Check if the name-email table matches expected data.

    Args:
        result: List of [name, email] rows from the getter
        expected: Expected data with 'authors' key containing list of [name, email] pairs
        **options: Additional options (e.g., case_sensitive)

    Returns:
        Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_authors = expected.get('authors', [])
    if not expected_authors:
        return 0.0
    result_sorted = sorted(result, key=lambda x: x[0].lower() if x and len(x) > 0 else '')
    expected_sorted = sorted(expected_authors, key=lambda x: x[0].lower() if x and len(x) > 0 else '')
    if len(result_sorted) != len(expected_sorted):
        return min(len(result_sorted) / len(expected_sorted), 1.0) * 0.3
    score = 0.0
    correct_rows = 0
    for (res_row, exp_row) in zip(result_sorted, expected_sorted):
        if len(res_row) < 2 or len(exp_row) < 2:
            continue
        name_match = res_row[0].strip().lower() == exp_row[0].strip().lower()
        email_match = res_row[1].strip().lower() == exp_row[1].strip().lower()
        if name_match and email_match:
            correct_rows += 1
        elif name_match:
            correct_rows += 0.5
    score = correct_rows / len(expected_sorted)
    return score

def check_email_not_in_recipients__c9ce3f52(result, expected, **options):
    """Check that specific emails are NOT in the recipient list AND required emails ARE included.

    This function performs two checks:
    1. Excluded emails (paid students) should NOT be in the recipient list
    2. Included emails (unpaid students) should ALL be in the recipient list

    Args:
        result: List of actual emails in To field
        expected: Dict with 'excluded_emails' and 'included_emails' keys
        **options: Additional options

    Returns:
        float: 1.0 if all conditions are met, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    excluded_emails = expected.get('excluded_emails', [])
    included_emails = expected.get('included_emails', [])
    result_normalized = [email.lower() for email in result]
    for excluded_email in excluded_emails:
        if excluded_email.lower() in result_normalized:
            return 0.0
    for included_email in included_emails:
        if included_email.lower() not in result_normalized:
            return 0.0
    return 1.0

def check_thunderbird_filter__947ae516(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_filter_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_filter_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_unique_senders__0399616e(result, expected, **options):
    """Compare unique sender list against expected values and verify Thunderbird state.

    Args:
        result: Dict with 'senders' (list), 'thunderbird_verified' (bool),
                'email_count_verified' (bool), or None if file doesn't exist
        expected: Dict with 'expected_senders' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
              - 0.0: File doesn't exist or completely wrong
              - 0.3: Correct output but Thunderbird verification failed
              - 0.5: Correct output, Thunderbird verified but email count not verified
              - 1.0: Full verification passed (output correct, Thunderbird verified, email count >= 5)
    """
    if result is None:
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    expected_senders = expected.get('expected_senders', [])
    if len(expected_senders) == 0:
        return 0.0
    senders = result.get('senders')
    if senders is None or not isinstance(senders, list):
        return 0.0
    if len(senders) == 0:
        return 0.0
    if len(senders) != len(expected_senders):
        return 0.0
    if senders != expected_senders:
        return 0.0
    thunderbird_verified = result.get('thunderbird_verified', False)
    email_count_verified = result.get('email_count_verified', False)
    if not thunderbird_verified:
        return 0.3
    if not email_count_verified:
        return 0.5
    return 1.0

def check_thunderbird_prefs__1f5ab547a29208fe81de884f80a1716d(result: str, expected: Dict[str, Dict[str, Dict[str, Any]]]) -> float:
    """
    Check Thunderbird preferences against expected rules.

    Args:
        result: Path to prefs.js file
        expected: Rules dict with 'expect' and optionally 'unexpect' keys

    Returns:
        1.0 if all expected preferences match and no unexpected preferences are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_rules = expected.get('expect', {})
    unexpect_rules = expected.get('unexpect', {})
    pref_pattern: Pattern[str] = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            match_: Match[str] = pref_pattern.match(line.strip())
            if match_ is None:
                continue
            key: str = match_.group('key')
            value = json.loads(match_.group('val'))
            if key in expect_rules:
                rule = expect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    expect_metrics[key] = value == ref
                elif method == 're':
                    expect_metrics[key] = bool(re.search(str(ref), str(value)))
                elif method == 'gt':
                    expect_metrics[key] = value > ref
                elif method == 'lt':
                    expect_metrics[key] = value < ref
                elif method == 'gte':
                    expect_metrics[key] = value >= ref
                elif method == 'lte':
                    expect_metrics[key] = value <= ref
            elif key in unexpect_rules:
                rule = unexpect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    unexpect_metric = unexpect_metric and value != ref
                elif method == 're':
                    unexpect_metric = unexpect_metric and (not bool(re.search(str(ref), str(value))))
    return float(all(expect_metrics.values()) and unexpect_metric)

def check_thunderbird_full_name__cd29a448(result: str, expected: dict) -> float:
    """
    Check if identity full name matches expected value.

    This function first identifies which identity ID corresponds to the email account
    'anonym-x2024@outlook.com' by checking mail.identity.id*.useremail preferences,
    then verifies that identity's fullName matches the expected value.

    Args:
        result: path to prefs.js file
        expected: dict with 'full_name' key

    Returns:
        float: 1.0 if name matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_name = expected.get('full_name')
    if not expected_name:
        logger.warning('No expected full_name provided')
        return 0.0
    target_email = 'anonym-x2024@outlook.com'
    try:
        identity_id = None
        with open(result, 'r') as f:
            for line in f:
                if 'mail.identity.id' in line and '.useremail' in line and ('user_pref' in line):
                    if target_email in line:
                        start = line.find('mail.identity.id')
                        if start != -1:
                            start += len('mail.identity.id')
                            end = line.find('.useremail', start)
                            if end != -1:
                                identity_id = 'id' + line[start:end]
                                logger.debug(f"Found identity ID '{identity_id}' for email '{target_email}'")
                                break
        if identity_id is None:
            logger.warning(f"Could not find identity ID for email '{target_email}'")
            return 0.0
        if not identity_id or not identity_id.startswith('id'):
            logger.warning(f"Invalid identity ID found: '{identity_id}'")
            return 0.0
        full_name_key = f'mail.identity.{identity_id}.fullName'
        with open(result, 'r') as f:
            for line in f:
                if full_name_key in line and 'user_pref' in line:
                    parts = line.split('"')
                    if len(parts) >= 4:
                        actual_name = parts[3]
                        if not actual_name or not actual_name.strip():
                            logger.warning(f"Full name is empty for identity '{identity_id}'")
                            return 0.0
                        logger.debug(f'Found fullName for {identity_id}: {actual_name}, expected: {expected_name}')
                        return 1.0 if actual_name == expected_name else 0.0
        logger.warning(f'{full_name_key} preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_reply_on_top__f50e981c(result: str, expected: dict) -> float:
    """
    Check if reply-on-top setting matches expected value.

    Args:
        result: path to prefs.js file
        expected: dict with 'reply_on_top' key (int: 0=bottom, 1=top)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_value = expected.get('reply_on_top')
    if expected_value is None:
        logger.warning('No expected reply_on_top value provided')
        return 0.0
    try:
        with open(result, 'r') as f:
            for line in f:
                if 'mail.identity.id1.reply_on_top' in line and 'user_pref' in line:
                    parts = line.split(', ')
                    if len(parts) >= 2:
                        value_part = parts[1].strip().rstrip(');')
                        actual_value = json.loads(value_part)
                        logger.debug(f'Found reply_on_top: {actual_value}, expected: {expected_value}')
                        return 1.0 if actual_value == expected_value else 0.0
        logger.warning('mail.identity.id1.reply_on_top preference not found')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_email_count_file__c95564e31fba18e60360c502646bdd5f(result: str, expected: dict, **options) -> float:
    """
    Check if a text file contains the correct email count.

    Args:
        result: Content from the created text file
        expected: Expected values from rules (dict with 'count')
        **options: Additional options

    Returns:
        1.0 if the content matches the expected count, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        result_str = str(result).strip()
        expected_count = expected.get('count', 1)
        if str(expected_count) in result_str:
            return 1.0
        try:
            count_value = int(result_str)
            if count_value == expected_count:
                return 1.0
        except ValueError:
            pass
        return 0.0
    except Exception:
        return 0.0

def check_tb_oldest3__dcd1e9d2ff68b2e8e650947cdbb3db16(result, expected, **options):
    """Compare extracted oldest 3 emails against expected values.

    Args:
        result: List of email dicts from getter
        expected: Dict with 'emails' key containing expected email list

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_emails = expected.get('emails', [])
    if len(result) != 3 or len(expected_emails) != 3:
        return 0.0
    try:
        result_dates = []
        for email in result:
            date_str = email.get('date', '')
            if isinstance(date_str, str):
                parsed_date = datetime.strptime(date_str, '%Y-%m-%d')
                result_dates.append(parsed_date)
            else:
                return 0.0
        for i in range(len(result_dates) - 1):
            if result_dates[i] >= result_dates[i + 1]:
                return 0.0
    except:
        return 0.0
    score = 0.0
    for (i, (res_email, exp_email)) in enumerate(zip(result, expected_emails)):
        if res_email.get('sender_name') == exp_email.get('sender_name') and res_email.get('sender_email') == exp_email.get('sender_email') and (res_email.get('date') == exp_email.get('date')):
            score += 1.0 / 3.0
    return score

def check_thunderbird_prefs__f0f94c7477cc672daff3cf5ec91ccac4(result: str, expected: Dict[str, Dict[str, Dict[str, Any]]]) -> float:
    """
    Check Thunderbird preferences against expected rules.

    Args:
        result: Path to prefs.js file
        expected: Rules dict with 'expect' and optionally 'unexpect' keys

    Returns:
        1.0 if all expected preferences match and no unexpected preferences are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_rules = expected.get('expect', {})
    unexpect_rules = expected.get('unexpect', {})
    pref_pattern: Pattern[str] = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            match_: Match[str] = pref_pattern.match(line.strip())
            if match_ is None:
                continue
            key: str = match_.group('key')
            value = json.loads(match_.group('val'))
            if key in expect_rules:
                rule = expect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    expect_metrics[key] = value == ref
                elif method == 're':
                    expect_metrics[key] = bool(re.search(str(ref), str(value)))
                elif method == 'gt':
                    expect_metrics[key] = value > ref
                elif method == 'lt':
                    expect_metrics[key] = value < ref
                elif method == 'gte':
                    expect_metrics[key] = value >= ref
                elif method == 'lte':
                    expect_metrics[key] = value <= ref
            elif key in unexpect_rules:
                rule = unexpect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    unexpect_metric = unexpect_metric and value != ref
                elif method == 're':
                    unexpect_metric = unexpect_metric and (not bool(re.search(str(ref), str(value))))
    return float(all(expect_metrics.values()) and unexpect_metric)

def check_test_emails_exported__a9ff16fc1be9331aa02283bf7c168b3e(result, expected, **options):
    """
    Check if the directory listing contains expected email files matching the test pattern.
    Verifies ONLY emails with 'Test' in the subject were exported (exclusivity check).

    Args:
        result: Path to file containing ls -R output
        expected: Expected patterns (dict with 'min_count' for minimum number of Test emails)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    min_count = expected.get('min_count', 1)
    test_pattern = re.compile('.*[Tt]est.*\\.eml')
    all_eml_pattern = re.compile('.*\\.eml')
    with open(result, 'r') as f:
        content = f.read()
        test_matches = test_pattern.findall(content)
        all_eml_matches = all_eml_pattern.findall(content)
    if len(test_matches) < min_count:
        return 0.0
    if len(all_eml_matches) != len(test_matches):
        return 0.0
    return 1.0

def check_thunderbird_two_folders__386310fb(result: dict, expected: dict, **options):
    """Check if two folders exist.

    Args:
        result: Dict from getter with personal_exists and work_exists
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_personal = expected.get('personal_exists', True)
    expected_work = expected.get('work_exists', True)
    score = 0.0
    if result.get('personal_exists') == expected_personal:
        score += 0.5
        logger.info('Personal folder check passed')
    else:
        logger.info(f"Personal folder check failed: got {result.get('personal_exists')}, expected {expected_personal}")
    if result.get('work_exists') == expected_work:
        score += 0.5
        logger.info('Work folder check passed')
    else:
        logger.info(f"Work folder check failed: got {result.get('work_exists')}, expected {expected_work}")
    return score

def check_thunderbird_attachment_complete__d38192b0(result_state, expected_state, **options):
    """
    Check if both Thunderbird compose window requirements are met.

    This metric verifies that:
    1. The Thunderbird compose window is still open (not closed or sent)
    2. The specified attachment is present in the compose window

    Both conditions must be true for the task to pass.

    Args:
        result_state: dict from getter with keys:
            - window_open: bool
            - attachment_present: bool
            - subject: str
            - attachment_name: str
        expected_state: dict with expected values (same structure)
        **options: Additional options

    Returns:
        float: 1.0 if both window is open and attachment is present, 0.0 otherwise
    """
    if not isinstance(result_state, dict):
        return 0.0
    window_open = result_state.get('window_open', False)
    if not window_open:
        return 0.0
    attachment_present = result_state.get('attachment_present', False)
    if not attachment_present:
        return 0.0
    return 1.0

def check_thunderbird_filter__07c1111d(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    DEFAULT_FOLDERS = ['Inbox', 'Drafts', 'Sent', 'Trash', 'Junk', 'Archive', 'Templates']
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            if _match_record(r, flt):
                if 'actionValue' in flt:
                    action_value = flt['actionValue']
                    is_custom = not any((default_folder in action_value for default_folder in DEFAULT_FOLDERS))
                    is_custom = is_custom and len(action_value) > 0
                    expect_metrics[i] = is_custom
                else:
                    expect_metrics[i] = False
            else:
                expect_metrics[i] = expect_metrics[i] or False
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_filter__782e51e3(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_filter__1da45b53(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_filter_record(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_filter_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_thunderbird_domain_count__1744c07b(result, expected, **options):
    """Compare domain count data against expected values.

    Args:
        result: Dict mapping domains to counts from getter
        expected: Dict with 'expected_counts' mapping
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_counts = expected.get('expected_counts', {})
    if len(result) != len(expected_counts):
        return 0.0
    matches = 0
    for (domain, count) in expected_counts.items():
        if result.get(domain) == count:
            matches += 1
    return matches / len(expected_counts) if expected_counts else 0.0

def check_thunderbird_filter__af973420(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check Thunderbird filter configuration.

    Args:
        result (str): path to filter def file
        expected (Dict[str, List[Dict[str, str]]]): dict like
          {
            "expect": [{key: value, condition: {...}}]
            "unexpect": [{key: value}]
          }

    Returns:
        float: 1.0 if filter matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    from desktop_env.evaluators.metrics.utils import _match_record
    filters: List[Dict[str, Union[str, List[str]]]] = []
    with open(result) as f:
        for l in f:
            if l.startswith('name='):
                filter_: Dict[str, Union[str, List[str]]] = {}
                filter_['name'] = _value_processor(l[6:-2])
            elif l.startswith('enabled='):
                filter_['enabled'] = _value_processor(l[9:-2])
            elif l.startswith('type='):
                filter_['type'] = _value_processor(l[6:-2])
            elif l.startswith('action='):
                filter_['action'] = _value_processor(l[8:-2])
            elif l.startswith('actionValue='):
                filter_['actionValue'] = _value_processor(l[13:-2])
            elif l.startswith('condition='):
                condition_str: str = _value_processor(l[11:-2])
                logger.debug('FILTER CONDITION: %s', condition_str)
                conditions: List[str] = _condition_pattern.findall(condition_str)
                logger.debug('FILTER CONDITIONS: %s', repr(conditions))
                filter_['condition'] = conditions
                logger.debug('FILTER %s', repr(filter_))
                filters.append(filter_)
    expect_metrics = [False] * len(expected.get('expect', []))
    unexpect_metric = True
    for flt in filters:
        for (i, r) in enumerate(expected.get('expect', [])):
            expect_metrics[i] = expect_metrics[i] or _match_filter_with_conditions(r, flt)
        unexpect_metric = unexpect_metric and (not any((_match_record(r, flt) for r in expected.get('unexpect', []))))
    return float(all(expect_metrics) and unexpect_metric)

def check_tb_email_summary__095c028a04cb1496af16d82dbb5f9c21(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the email summary statistics match expected values.

    Args:
        result: Dict from getter containing summary statistics
        expected: Dict with expected summary values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Invalid result format')
        return 0.0
    score = 0.0
    checks = 0
    if 'total_count' in expected:
        checks += 1
        if result.get('total_count') == expected['total_count']:
            score += 1.0 / 3
        else:
            logger.debug(f"Total count mismatch: got {result.get('total_count')}, expected {expected['total_count']}")
    if 'unique_senders' in expected:
        checks += 1
        if result.get('unique_senders') == expected['unique_senders']:
            score += 1.0 / 3
        else:
            logger.debug(f"Unique senders mismatch: got {result.get('unique_senders')}, expected {expected['unique_senders']}")
    if 'total_attachments' in expected:
        checks += 1
        if result.get('total_attachments') == expected['total_attachments']:
            score += 1.0 / 3
        else:
            logger.debug(f"Total attachments mismatch: got {result.get('total_attachments')}, expected {expected['total_attachments']}")
    return score if checks > 0 else 0.0

def check_tb_enotices__5ef2369a29068457d8de24f5079f30bb(result, expected, **options):
    """Compare extracted eNotices emails against expected values.

    Args:
        result: List of email dicts from getter
        expected: Dict with 'emails' key containing expected email list

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_emails = expected.get('emails', [])
    if len(result) != len(expected_emails):
        return 0.0
    score = 0.0
    for (res_email, exp_email) in zip(result, expected_emails):
        matches = 0
        total = 4
        if res_email.get('sender_name') == exp_email.get('sender_name'):
            matches += 1
        if res_email.get('sender_email') == exp_email.get('sender_email'):
            matches += 1
        if res_email.get('subject') == exp_email.get('subject'):
            matches += 1
        if res_email.get('cc') == exp_email.get('cc'):
            matches += 1
        score += matches / total / len(expected_emails)
    return score

def check_thunderbird_account_name__dced30f9(result: str, expected: Dict[str, Any]) -> float:
    """
    Check if the account display name in Thunderbird prefs matches the expected value.

    This function:
    1. Searches for the server ID associated with 'anonym-x2024@outlook.com' by checking
       mail.server.serverN.userName preferences
    2. Once found, checks if that server's mail.server.serverN.name equals 'Work Email'

    Args:
        result: path to prefs.js file
        expected: dict with 'server_name' key containing expected display name

    Returns:
        float: 1.0 if name matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_name = expected.get('server_name')
    if not expected_name:
        logger.warning('No expected server_name provided')
        return 0.0
    try:
        target_email = 'anonym-x2024@outlook.com'
        server_id = None
        with open(result, 'r') as f:
            content = f.read()
        for line in content.split('\n'):
            if 'mail.server.server' in line and '.userName' in line and ('user_pref' in line):
                parts = line.split('"')
                if len(parts) >= 4:
                    pref_name = parts[1]
                    pref_value = parts[3]
                    if pref_value == target_email:
                        server_id = pref_name.split('.')[2]
                        logger.debug(f"Found email account '{target_email}' associated with {server_id}")
                        break
        if not server_id:
            logger.warning(f"Could not find server ID for email account '{target_email}' in prefs.js")
            return 0.0
        pref_to_find = f'mail.server.{server_id}.name'
        for line in content.split('\n'):
            if pref_to_find in line and 'user_pref' in line:
                parts = line.split('"')
                if len(parts) >= 4:
                    actual_name = parts[3]
                    logger.debug(f'Found server name: {actual_name}, expected: {expected_name}')
                    return 1.0 if actual_name == expected_name else 0.0
        logger.warning(f'{pref_to_find} preference not found in prefs.js')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading prefs.js: {e}')
        return 0.0

def check_thunderbird_prefs__4d8c063289348a4e74b61401db33c326(result: str, expected: Dict[str, Dict[str, Dict[str, Any]]]) -> float:
    """
    Check Thunderbird preferences against expected rules.

    Args:
        result: Path to prefs.js file
        expected: Rules dict with 'expect' and optionally 'unexpect' keys

    Returns:
        1.0 if all expected preferences match and no unexpected preferences are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_rules = expected.get('expect', {})
    unexpect_rules = expected.get('unexpect', {})
    pref_pattern: Pattern[str] = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            match_: Match[str] = pref_pattern.match(line.strip())
            if match_ is None:
                continue
            key: str = match_.group('key')
            value = json.loads(match_.group('val'))
            if key in expect_rules:
                rule = expect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    expect_metrics[key] = value == ref
                elif method == 're':
                    expect_metrics[key] = bool(re.search(str(ref), str(value)))
                elif method == 'gt':
                    expect_metrics[key] = value > ref
                elif method == 'lt':
                    expect_metrics[key] = value < ref
                elif method == 'gte':
                    expect_metrics[key] = value >= ref
                elif method == 'lte':
                    expect_metrics[key] = value <= ref
            elif key in unexpect_rules:
                rule = unexpect_rules[key]
                method = rule.get('method', 'eq')
                ref = rule.get('ref')
                if method == 'eq':
                    unexpect_metric = unexpect_metric and value != ref
                elif method == 're':
                    unexpect_metric = unexpect_metric and (not bool(re.search(str(ref), str(value))))
    return float(all(expect_metrics.values()) and unexpect_metric)

def check_thunderbird_all_emails_read__a63fb94e(result: dict, expected: dict, **options):
    """Check if all emails are marked as read.

    Args:
        result: Dict from getter with all_read status
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_all_read = expected.get('all_read', True)
    if result.get('all_read') == expected_all_read:
        logger.info(f"All emails read check passed: {result.get('read_emails')}/{result.get('total_emails')} emails are read")
        return 1.0
    else:
        logger.info(f"All emails read check failed: {result.get('read_emails')}/{result.get('total_emails')} emails are read, expected all_read={expected_all_read}")
        return 0.0

def check_thunderbird_draft_attachment__d38192b0_aug9(result_state, expected_state, **options):
    """
    Check that the Thunderbird email draft has the correct attachment and is not sent.

    This metric verifies:
    1. Email exists in Drafts folder
    2. Email is NOT in Sent folder (do not send requirement)
    3. Email has the expected attachment
    4. Attachment is not empty (has reasonable size)

    Args:
        result_state (dict): Result from get_thunderbird_draft_attachment getter
            {
                'draft_exists': bool,
                'not_sent': bool,
                'has_attachment': bool,
                'attachment_name': str or None,
                'attachment_size': int or None
            }
        expected_state (dict): Expected state
            {
                'draft_exists': bool (default True),
                'not_sent': bool (default True),
                'has_attachment': bool (default True),
                'min_attachment_size': int (default 500)
            }
        **options: Additional options
            - partial_credit: bool (default True) - Give partial credit for each requirement

    Returns:
        float: Score (0.0 to 1.0)
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    expected_draft_exists = expected_state.get('draft_exists', True)
    expected_not_sent = expected_state.get('not_sent', True)
    expected_has_attachment = expected_state.get('has_attachment', True)
    min_attachment_size = expected_state.get('min_attachment_size', 500)
    partial_credit = options.get('partial_credit', True)
    draft_exists = result_state.get('draft_exists', False)
    not_sent = result_state.get('not_sent', False)
    has_attachment = result_state.get('has_attachment', False)
    attachment_size = result_state.get('attachment_size', 0)
    checks = {'draft_exists': draft_exists == expected_draft_exists, 'not_sent': not_sent == expected_not_sent, 'has_attachment': has_attachment == expected_has_attachment, 'attachment_size_ok': attachment_size is not None and attachment_size >= min_attachment_size if has_attachment else True}
    logger.info(f"Draft exists: {draft_exists} (expected: {expected_draft_exists}) - {('PASS' if checks['draft_exists'] else 'FAIL')}")
    logger.info(f"Not sent: {not_sent} (expected: {expected_not_sent}) - {('PASS' if checks['not_sent'] else 'FAIL')}")
    logger.info(f"Has attachment: {has_attachment} (expected: {expected_has_attachment}) - {('PASS' if checks['has_attachment'] else 'FAIL')}")
    logger.info(f"Attachment size: {attachment_size} (min: {min_attachment_size}) - {('PASS' if checks['attachment_size_ok'] else 'FAIL')}")
    if partial_credit:
        score = 0.0
        score += 0.25 if checks['draft_exists'] else 0.0
        score += 0.25 if checks['not_sent'] else 0.0
        score += 0.25 if checks['has_attachment'] else 0.0
        score += 0.25 if checks['attachment_size_ok'] else 0.0
        logger.info(f'Partial credit score: {score}')
        return score
    elif all(checks.values()):
        logger.info('All checks passed - score: 1.0')
        return 1.0
    else:
        logger.warning(f'Some checks failed - score: 0.0')
        return 0.0

def check_thunderbird_socket_timeout__a3ecd62701cd6d8575ca51c05b0fdb89(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if Thunderbird server socket and timeout settings are configured correctly.

    Args:
        result: Path to prefs.js file
        expected: Expected configuration with keys:
            - socket_type: expected int value
            - timeout: expected int value
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    pref_pattern = re.compile('^user_pref\\("(?P<key>(?:[^"]|\\\\")+)\\", (?P<val>.+)\\);$')
    found_socket = False
    found_timeout = False
    socket_matches = False
    timeout_matches = False
    with open(result, 'r') as f:
        for line in f:
            match = pref_pattern.match(line.strip())
            if match is None:
                continue
            key = match.group('key')
            value = json.loads(match.group('val'))
            if key == 'mail.server.server1.socketType':
                found_socket = True
                socket_matches = value == expected.get('socket_type', 3)
            elif key == 'mail.server.server1.timeout':
                found_timeout = True
                timeout_matches = value == expected.get('timeout', 29)
    score = 0.0
    if found_socket and socket_matches:
        score += 0.5
    if found_timeout and timeout_matches:
        score += 0.5
    return score

def check_thunderbird_bcc__497f2c49(result: dict, expected: dict) -> float:
    """
    Check if BCC setting matches expected value for the specific email account.

    Args:
        result: dict containing identity_id, do_bcc, bcc_list, and email
        expected: dict with 'do_bcc' key (bool: true/false)

    Returns:
        float: 1.0 if all BCC settings are correctly configured, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_value = expected.get('do_bcc')
    if expected_value is None:
        logger.warning('No expected do_bcc value provided')
        return 0.0
    if result.get('identity_id') is None:
        logger.warning(f"Could not find identity ID for email account {result.get('email')}")
        return 0.0
    actual_do_bcc = result.get('do_bcc')
    if actual_do_bcc is None:
        logger.warning(f"doBcc setting not found for identity {result.get('identity_id')}")
        return 0.0
    if actual_do_bcc != expected_value:
        logger.debug(f'doBcc mismatch: actual={actual_do_bcc}, expected={expected_value}')
        return 0.0
    if expected_value is True:
        bcc_list = result.get('bcc_list')
        if bcc_list is None or bcc_list == '':
            logger.warning(f'BCC is enabled but no BCC address is configured (doBccList is empty or missing)')
    logger.debug(f"BCC settings verified for identity {result.get('identity_id')}: do_bcc={actual_do_bcc}, bcc_list={result.get('bcc_list')}")
    return 1.0
