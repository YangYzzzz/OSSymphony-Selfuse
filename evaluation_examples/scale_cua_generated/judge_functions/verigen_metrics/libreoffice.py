"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_libreoffice_ext_mgr__441b7b9ef584e98e72cbacfa6a30a033', 'check_libreoffice_dict__17aed94a4f09fd9062bfe0ffbfb8f904', 'check_all_headoffice_officers__beab95b2414a2d1c1b2b4df1d8d7954b', 'check_headoffice_officer__ab6378f137ebc3c17ec9a3b004bb1bd0', 'check_libreoffice_template__5041b595bf111dbc12a85f19ac3493af', 'check_libreoffice_command__14fa116b50d2b7c945d2ede19b12f134', 'check_unique_officer_count__727a19e2a3d544ebb2bea225e698705f', 'check_officer_match__c0b425ace9ff45f341465bcadf538b3f']

def check_libreoffice_ext_mgr__441b7b9ef584e98e72cbacfa6a30a033(result: str, expected: dict, **options) -> float:
    """
    Check if expected extension patterns are found

    Args:
        result: Output from command
        expected: Dictionary with 'expect' patterns to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_patterns: List[Pattern[str]] = [re.compile(ptt) for ptt in expected.get('expect', [])]
    if not expect_patterns:
        return 0.0
    for pattern in expect_patterns:
        if not pattern.search(result):
            logger.info(f'Pattern not found: {pattern.pattern}')
            return 0.0
    logger.info('All expected patterns found')
    return 1.0

def check_libreoffice_dict__17aed94a4f09fd9062bfe0ffbfb8f904(result: str, expected: dict, **options) -> float:
    """
    Check if expected dictionary patterns are found in LibreOffice directory

    Args:
        result: Output from grep command
        expected: Dictionary with 'expect' patterns to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_patterns: List[Pattern[str]] = [re.compile(ptt) for ptt in expected.get('expect', [])]
    if not expect_patterns:
        return 0.0
    for pattern in expect_patterns:
        if not pattern.search(result):
            logger.info(f'Pattern not found: {pattern.pattern}')
            return 0.0
    logger.info('All expected patterns found')
    return 1.0

def check_all_headoffice_officers__beab95b2414a2d1c1b2b4df1d8d7954b(result: List[Dict[str, str]], expected: Dict[str, Any], **options) -> float:
    """Check if all rows for a HeadOffice have the correct officer assigned.

    Args:
        result: List of dicts with 'area' and 'officer' from getter
        expected: Dict with 'expected_officer' and 'expected_count'
        **options: Additional options

    Returns:
        Score based on correct officer assignments (partial credit)
    """
    if not isinstance(result, list) or not isinstance(expected, dict):
        return 0.0
    expected_officer = expected.get('expected_officer', '')
    expected_count = expected.get('expected_count', 0)
    if len(result) != expected_count:
        return 0.0
    if not result:
        return 0.0
    correct_count = 0
    for entry in result:
        if entry.get('officer', '').strip() == expected_officer.strip():
            correct_count += 1
    return correct_count / len(result) if result else 0.0

def check_headoffice_officer__ab6378f137ebc3c17ec9a3b004bb1bd0(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if Officer Name column is populated AND the most frequent HeadOffice has correct officer assigned.

    This metric verifies BOTH requirements from the instruction:
    1. Officer Name column (F2:F12) is populated with VLOOKUP results
    2. Most frequent HeadOffice is correctly identified with its officer

    Args:
        result: Dict from getter with:
            - 'headoffice': Most frequent HeadOffice
            - 'officers': Officers for that HeadOffice
            - 'all_officers_populated': Whether all Officer Name cells are filled
            - 'populated_count': Number of filled cells
            - 'total_rows': Total rows expected (11)
            - 'vlookup_samples_correct': Whether sample VLOOKUP results are correct
        expected: Dict with 'headoffice' and 'expected_officer'
        **options: Additional options

    Returns:
        Score breakdown:
        - 0.0: Failed both requirements
        - 0.3: Column populated but headoffice/officer wrong
        - 0.4: Headoffice correct but column not populated
        - 0.6: Headoffice correct AND column populated
        - 0.8: Headoffice + officer correct but column not fully populated (>= 80% filled)
        - 1.0: All requirements met (column populated + headoffice + officer correct)
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_headoffice = expected.get('headoffice', '')
    expected_officer = expected.get('expected_officer', '')
    actual_headoffice = result.get('headoffice', '')
    actual_officers = result.get('officers', [])
    all_officers_populated = result.get('all_officers_populated', False)
    populated_count = result.get('populated_count', 0)
    total_rows = result.get('total_rows', 11)
    vlookup_samples_correct = result.get('vlookup_samples_correct', False)
    population_ratio = populated_count / total_rows if total_rows > 0 else 0
    column_populated = all_officers_populated or (population_ratio >= 0.8 and vlookup_samples_correct)
    headoffice_correct = str(actual_headoffice).strip() == str(expected_headoffice).strip()
    officer_correct = expected_officer in [str(o).strip() for o in actual_officers]
    score = 0.0
    if column_populated and (not headoffice_correct):
        score = 0.3
    elif headoffice_correct and (not column_populated):
        score = 0.4
    elif headoffice_correct and column_populated and (not officer_correct):
        score = 0.6
    elif headoffice_correct and officer_correct and (not column_populated):
        score = 0.8
    elif headoffice_correct and officer_correct and column_populated:
        score = 1.0
    return score

def check_libreoffice_template__5041b595bf111dbc12a85f19ac3493af(result: str, expected: dict, **options) -> float:
    """
    Check if expected template patterns are found

    Args:
        result: Output from command
        expected: Dictionary with 'expect' patterns to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_patterns: List[Pattern[str]] = [re.compile(ptt) for ptt in expected.get('expect', [])]
    if not expect_patterns:
        return 0.0
    for pattern in expect_patterns:
        if not pattern.search(result):
            logger.info(f'Pattern not found: {pattern.pattern}')
            return 0.0
    logger.info('All expected patterns found')
    return 1.0

def check_libreoffice_command__14fa116b50d2b7c945d2ede19b12f134(result, expected, **options):
    """Check if bash history contains libreoffice conversion command for .doc files.

    Args:
        result: Bash history content
        expected: Expected pattern rules
        **options: Additional options

    Returns:
        float: 1.0 if correct command found, 0.0 otherwise
    """
    pattern = '(soffice|libreoffice).+--convert-to\\s+pdf.+\\*\\.doc(?:\\s|$)'
    if re.search(pattern, result):
        return 1.0
    else:
        return 0.0

def check_unique_officer_count__727a19e2a3d544ebb2bea225e698705f(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if VLOOKUP is correct and unique officer count matches expected.

    This verifies TWO requirements:
    1. The VLOOKUP mappings are correct (each officer in column F matches the expected officer for the branch in column E)
    2. The unique officer count matches the expected value

    Args:
        result: Dict from getter with:
            - 'unique_count': Number of unique officers found
            - 'officers': List of unique officer names
            - 'vlookup_correct': Boolean indicating if VLOOKUP mappings are correct
        expected: Dict with:
            - 'unique_count': Expected number of unique officers
        **options: Additional options

    Returns:
        1.0 if VLOOKUP is correct AND count matches, 0.0 otherwise
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    vlookup_correct = result.get('vlookup_correct', False)
    if not vlookup_correct:
        return 0.0
    expected_count = expected.get('unique_count', 0)
    actual_count = result.get('unique_count', 0)
    if actual_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_officer_match__c0b425ace9ff45f341465bcadf538b3f(result: Dict[str, Dict[str, str]], expected: dict, **options) -> float:
    """Check if officer names for all areas match expected values.

    Args:
        result: Dict mapping area names to their data (from getter)
                e.g., {"Vastral": {"head_office": "Ahmedabad", "officer": "C. D. Dey"}, ...}
        expected: Dict with 'area_officers' key containing expected area-officer mappings
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 (partial credit based on correct count / total count)
    """
    if not isinstance(expected, dict):
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    expected_mappings = expected.get('area_officers', {})
    if not expected_mappings:
        return 0.0
    total_areas = len(expected_mappings)
    correct_count = 0
    for (area_name, expected_data) in expected_mappings.items():
        if area_name not in result:
            continue
        result_data = result[area_name]
        expected_officer = expected_data.get('officer', '')
        result_officer = result_data.get('officer', '')
        if str(result_officer).strip() == str(expected_officer).strip():
            correct_count += 1
    if total_areas == 0:
        return 0.0
    return correct_count / total_areas
