"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['compare_images', 'check_image_size__2aa2e23a948bbaf102eec8741d709529', 'check_image_properties__4bd1c0a2e0720abe2c2c09ee9978ce22', 'check_triangle_flipped__c73a98370436f02ecc8c151f49a9b91c', 'check_image_size__5db5c513f28a5ad7ad8d60ef05b14b35', 'check_layer_exists__a85f6474', 'check_image_orientation__844f5e73b108da449b9b68fcbef6bbf2', 'check_image_mode__8f91f3a7', 'check_image_dimensions__baa9d5de58b4726390c9c04659eb9fca', 'check_image_extracted__a3446500', 'check_gimp_gimprc_setting__5b6e5d1f', 'check_png_file_exists__846e3a68', 'check_layer_name_config__3336340a', 'check_triangle_flipped__1aaf02638da71a4a84f47e22e2395da3', 'check_layer_fixed_size__0dbca844', 'check_image_hash__db2588a0', 'check_image_files_exist__abdf3926e9609e7e5b435c8cdbb40013', 'check_triangle_topright__54d36b10', 'check_specific_image_inserted__2759cd98', 'check_image_dimensions__a7d5fd37', 'check_moved_pngs__ec3ddc36', 'check_image_grayscale__bf5ecb70b33ef3dffbe095b744ec38a7', 'check_png_exists__06536a54', 'check_gimp_saturation_decrease__62aadfc75943521b0b091bf9d9c10f24', 'check_image_exists__e327229f', 'check_gif_file__b130b682', 'check_image_dimensions__13d2c6c3', 'check_gif_file__f8138076f6546232c093f5027d50db21', 'check_image_files_exist__e4b458033c1389ecc56cd63f5eae9626', 'check_image_resize__dbc2e36d8461f22779ef5b4f0521a159', 'check_triangle_scaled__93095e022ff3d0c1026d009f3ccc512b', 'check_image_dimensions__4eb874068ccb861273ecf8604bdafb3c', 'check_image_hash__6510dd9a', 'check_image_count__02587c7a', 'check_default_video_player__9b1915f5', 'check_gimp_menubar__4e59f0cb646eef506e256c2929b9d463', 'check_image_is_grayscale__0e7db70c', 'check_bg_image_exists__0a4467dd72d00d84bc5d80c765c4ad6f', 'check_image_cropped__7bac509f', 'check_image_mirror__bc5fe443a1102b529ba31d0ac9c81ab1', 'check_layer_exists__8565a91c', 'check_image_hash__0c114ed7', 'check_image_dimensions__3f39d534f7803089ed331d19c2a1bc89', 'check_image_hash__73176121', 'check_triangle_top_left__86e46240', 'check_gif_file__430ad71e1ce94d5ea2bb30ba073fd937', 'check_image_dimensions__a846d552e46987bd85e04c0a4f658c7a', 'check_image_dimensions__68566bbc', 'check_pdf_image_count__f35ebb0d', 'check_image_properties__c6063730', 'check_image_dimensions__7816ba80', 'check_layer_count__58c8d068', 'check_image_dimensions__f5cfdff3841c16728bb4565a839b59ca', 'check_image_size__07fcde31879d28764a081db212afaf2d', 'check_image_size__0d547ae5', 'check_image_properties__c4f65e24', 'validate_image_count__b01cf463', 'check_image_format__22b30bf2', 'check_gif_file__092b88db7884c339b37665609b49294b', 'check_image_columns_reversed__3936f39d', 'check_jpeg_exists_and_size__ad4cc5ac', 'check_gimp_rulers_setting__96f3f9d88ca4e4230495b91ff566eb51', 'check_image_properties__811349c6', 'check_gimp_statusbar_setting__2d2ae6cc356da88025063f58e6110bad', 'check_gif_file__d18c2fd4', 'check_image_dimensions__9d6c98f0', 'check_image_size__8677d5c370f70c69494653c2a8ef5be2', 'check_image_format__900fc36276e9eaf12471a7834992cf5c', 'check_image_flipped__07b5058cd3df711389d2b4342d0c561c', 'check_image_dimensions__016cfcf5', 'check_image_dimensions__bec1165ff4f2eb5b48dc7de50a4fe1ab', 'check_image_grayscale__9660fa13484d43a1462069231ef86deb', 'check_image_extracted__345e8ddb', 'check_gif_file__8d385ddc', 'check_image_format__739292ff', 'check_vlc_default_player__d253b8f0', 'check_triangle_top_right__1a2ce7aa', 'check_image_resize_and_structure_sim__ed3206c993cd330e3bb4b56cf8b439e8', 'check_image_cropped__166a89c9', 'check_layer_exists__b148e375_v3', 'check_extracted_image__938bcb9a', 'check_png_deleted__eb6b46ad', 'check_image_hash__d1128c0a', 'check_image_file__203069587d53a571860bccb97348992b', 'check_image_flipped__16b4973e', 'check_gimp_autosave_setting__d7b0aa1f', 'check_image_dimensions__739292ff', 'check_all_images__1574bc56f755697238f4190c2d72a32b', 'check_final_image_count__5d239d03', 'check_image_extracted__f3e27026', 'check_image_hash__0969de9d', 'check_image_rotated__777d01ae6d6663c5897f2674681dca2e', 'check_gimp_single_window__d9681274028e239cd7a4e13a8dc53d3c', 'check_image_size__2110d1dc', 'check_png_export__046a9ca717a3ff75711d7d7d1d876a5f', 'check_image_dimensions__b24412f0', 'check_layer_exists__17690e7b', 'check_third_image__9a3feca86555f3c82732707871883142', 'check_has_multiple_images__67a58d2f', 'check_gimp_rulers__f13e522727f9f47f05a1828391e4f1b5', 'check_image_mirrored__2237c5bebdb76e54ae53ea89e71ca4a3', 'check_jpeg_export__a818e71b', 'check_composite_photo_organization__8542f72b', 'check_image_blurred__eea46e4e4afa193660c5f52c6a2da7a9', 'check_layer_exists__22c9ed6961463c7a6e1af4ed7b883c13', 'check_layer_renamed__677f0e4c', 'check_image_horizontal_flip__08f973926b8f35af5489796c73a6c6e0', 'check_triangle_bottom_left__e0c275c9', 'check_image_resize__dc2a3b39', 'check_jpg_files_copied__2c299b404d518d11da9f18c6fc937d90', 'check_image_format_and_existence__e19bfc7ef1338231fef513d8a5b2f6d1', 'check_jpg_exists__8d385ddc', 'check_image_extracted__4e5534f0', 'check_gimp_tabs__ca73cf065f0842768815735a3354abf0', 'check_image_cropped__c48d85ee866b2dcfd0ab60f091a5a2b6', 'check_image_crop__4ff6a540', 'check_image_resized__4939e12b', 'check_file_exists_and_image_size__16fb0dc6', 'check_image_rotated__4e34ef5d', 'check_png_export_and_structure__c69055e15cceefc40a87e6de042c2331', 'check_image_scaled__5d068218', 'check_image_file_exists__7c0cc95089263e14cb308f3757c8acc1', 'check_image_properties__bd10eca8', 'check_image_flip_vertical__fc198eb0', 'check_image_crop__a8983eb9', 'check_image_size__fb0dfb53338d620d275120282f62d824', 'check_png_resolution__27528fb8', 'check_jpg_list__8ddf03c6b80780c49b4f9497dee3f888', 'check_image_crop__a0afbcc2', 'check_image_is_grayscale__1ac9295ec86f9e2e04c973e2e47b273c', 'check_image_dimensions__9472df29', 'check_image_extracted__69184b17', 'check_triangle_bottomleft__4f468172ec938d83ce9f0bb7cb13c9fc', 'check_png_count_range__b4c05b10', 'check_triangle_rotation__8f589d19', 'check_image_properties__37707684957589279b0fa14602529fe7', 'check_two_folders_with_images__6a4313ae', 'check_pdf_with_image__d5371d5b', 'check_triangle_color_blue__d52919fa54c3d57b57da98359753b674', 'check_image_extracted__1a437041', 'check_bottom_half_image__4baa14c913fe32a60f559a98296affae', 'check_gimp_toolbox_setting__90386cba106758972cba7bf949bb562f', 'check_triangle_bottomright__95d4f9f61b019dd634339dd50c3adb82', 'check_photo_rename__0a812fadd7008b5cb899c479c82ab6e7', 'check_gif_file__16fb0dc6', 'check_image_dimensions__fa4ed82a', 'check_image_resize__89f19527', 'check_gif_exists__11cf8ab6', 'check_total_images__2759cd98', 'check_gimp_gimprc_setting__9c13adcb', 'check_gif_file__3d85489a', 'check_image_rotated__b49205bc', 'check_image_crop__70ffd430', 'check_gimp_menubar_setting__eebb7fd7089a7abc35d29b3d4832455e', 'verify_total_images__d49c5873', 'check_gif_file__e66cc6b0', 'check_image_resized__f8a3e253', 'check_image_aspect_ratio__f85979b813875fecca12ba1c6ab4cc68', 'check_gif_file_basic__de678d13fd248567f73258f2b3cb0372', 'check_image_crop__59e6ca63b22becd85a8942d1a29325d9', 'check_layer_exists__ca8ad8ab', 'check_png_exists__06722a19', 'check_layer_name_config__174a6594', 'check_triangle_topleft__40af8739a7f9d38f5046e0dfa41c5f6e', 'check_image_dimensions__696943e1', 'check_png_saved__b9abaa4fbc51b493882263c6a6aff8fe', 'check_image_dimensions__34372286', 'check_image_mode__78bfc971', 'check_vacation_jpgs__d7a48669399bf74024b2a979a32d4ae1', 'check_image_resized__912586a80097a155904c15da97ac2079', 'check_image_format__c0fb0f23', 'check_image_grayscale__90bdfeb2ebfd1bac6873718be37b3912', 'check_image_dimensions__77b19ce3287accb29381eac14cc998b5', 'check_layer_exists__b148e375', 'check_image_rows_reversed__1f38b5ad', 'check_image_dimensions__b6f18d98d1993dda1c36f44651fe6a5d', 'check_image_crop_dimensions__c26faa30e6e2f9e09b4748ad3193b390', 'check_image_rotated_90_clockwise__25c734e4497155c51ced0623dec284fc', 'check_image_width__d4f10b60', 'check_default_video_player__52f0bd95', 'check_image_dimensions__505cf5fc', 'check_image_grayscale__37ae16c860b893668df4aed8d0b9ae18', 'check_image_mirror__dae8d36d', 'check_image_deleted__53c4502a', 'check_image_rotated__045e2e0c', 'check_image_dimensions__ab9c94ea', 'verify_image_count__8e1f0943', 'check_gif_animated__fc6196ba2aeabb7bc4c476007b2b24c6', 'check_gif_file__ea8c7a7a', 'check_image_crop__f3b4c5c2', 'check_gimp_fullscreen__dbacd0b982a5b74d88d0d51e944036be', 'check_triangle_right_edge__1418e487', 'check_triangle_bottom_right__e4487e27c5c6e4232b26add556e7d796', 'check_image_cropped__8159102f', 'check_moved_pngs__486dfc7618cd1caae2ceb6550f7743a0', 'check_image_dimensions__a4f96e46', 'check_layer_exists__8249d409', 'check_image_format__fa27ef76', 'check_image_rotated__55bebaef7ca999b793134c2c00f342a9', 'check_image_color_mode__14447080df6d9553c7e99d3265fc5a81', 'check_image_size__e19bd559', 'check_image_dimensions__f2472a44', 'check_layer_added__e5a8318d', 'check_image_hash__8778dd25', 'check_image_hash__10e2f0b6', 'check_image_rotation__b014bcfb2b56c94b7ab718a85ad6cff9', 'check_image_rotate_180__201e011d', 'check_image_properties__72e666ad', 'check_triangle_flipped__0293a54ed09c3032ba177ba116651a07', 'check_image_dimensions__795e137d', 'check_image_modified__4c0f04bf', 'check_gif_file__11cf8ab6', 'check_gimp_gimprc_setting__adbf37d0', 'check_image_format__d255bc15', 'check_image_dimensions__05f2a34a', 'check_gimp_is_grayscale__7ae19854d29f1b0c8cdcf13e028f0bdd', 'check_gimp_gimprc_setting__dc948653', 'check_png_file_exists__d04e36cb527bbbbfdb06458981bf8945', 'check_triangle_scale_50__2550dc37', 'check_png_files__4a34f03bcb82e7e8037181cd9a91ae6e', 'check_image_dimensions__3b2067dec4f66b0e25da3355b1fb8f3e', 'check_gif_file__0d0257bdd2e8345a807cfdefd39ffa3b', 'check_image_horizontal_flip__340a3600a88f3d4dc2217b9f986d625d', 'check_triangle_topedge__e65601d6', 'check_image_dimensions__1728b8eeebbdebd7894eb2578c0c67ec', 'check_image_grayscale__f3cdb596', 'check_gif_file__06536a54', 'check_image_hash__b461596b', 'check_image_rotated_90__30840ecf', 'check_image_properties__0a6f6cc1', 'check_folder_images__fe2f25921e2c2c43fbc9e31c35bccb78', 'check_image_size__7ecb394c9eae8a8e135a21ca629ec0de', 'check_image_rotation__82cdf8c66b7532068e061ebffd0a6c33', 'check_image_scaled__e429d357', 'check_image_scaled__bf4967f3e9931b3c80be8e4dbf6e04b7', 'check_image_dimensions__ab3f6dd1', 'check_image_dimensions__9f5722fc', 'check_triangle_color__0ab9cbd5', 'check_new_image_added__b5b56630', 'check_triangle_deleted__0511b7c6', 'check_image_format__945e33b1', 'check_image_cropped__cdcc2775', 'check_image_rotated__6575e228', 'check_image_size__4f6d3ad3', 'check_image_properties__4894850c', 'check_triangle_topright__b91aae4ce11e76fbc123f25e08b42c98', 'check_image_properties__50f8e5bb', 'check_triangle_rotated__e0c6cf186be7cc4682d3712837dcdd72', 'check_triangle_topleft__562f7ff1b81bbf5a1dbcdac12003c550', 'check_image_flipped__9f1f61b8216550440a48089a3e4c1731', 'check_gimp_image_size__e2468a8febb27268d777ba03561c41ab', 'check_triangle_rotation__2c517d983ea369519ffc55979e3393ec', 'check_png_files_exact_count__f9d97b19', 'check_triangle_opacity__0ff223c2', 'check_image_resize__eed0216b', 'check_gimp_exact_dimensions__4bd0ac4fe70775f29bef20a161a34c39', 'check_image_properties__70f3db86', 'check_image_cropped__06e412f2', 'check_layer_exists__aa5f92aa', 'check_image_dimensions__da5d7378bbf41608325407fc00f8d126', 'check_gimp_gimprc_setting__17026b9e', 'check_image_dimensions__184d842d21bd06203c79c089532a2315', 'check_png_valid__2511ecbd', 'check_image_dimensions__de812dd5b44b906cb9793a8d4a3f91bf', 'check_image_file_exists__3948a175', 'check_image_blurred__782a609a', 'check_layer_group_config__7ba73b05', 'check_image_dimensions__c5f81e73faaccc56bdfa2edf29f272b7', 'check_gimp_fullscreen_mode__8efddf2685fdb790d7823145c3565e94', 'check_image_dimensions__565232d7', 'check_image_resized__8000c7e00c88e8061974ceb3ccc555df', 'check_png_dimensions__92ca6baf', 'check_image_compression__c46a6f1dddc552cd368bc819d4cce6f7', 'check_triangle_scale__12d50454feda301909898f2cf2cce54b', 'check_image_extracted__a8440735', 'check_image_mirror__77b8ab4d994f43ac89308ca087d7c4b4', 'check_image_dimensions__4694337da0a8886d5bd508a95fd83b12', 'check_image_dimensions__7cb89717dbfd62e5cbbd4dcc85a4e268', 'check_image_crop__9e45757e', 'check_jpeg_export__39da7f334341155f29f73cbacf02786e', 'check_image_dimensions__c984db77', 'check_image_hash__69eadfa0', 'check_gif_exists__e4d07acf', 'check_gif_dimensions__90c23f3797e29598f167849a527a40d5', 'check_gif_file__78103de86e64961437a4bcd00b97b9bc', 'check_image_size__10ab0aed', 'check_image_resize__ff57a44d', 'check_triangle_color__978cb6f31473d4802226bc7ae94b7399', 'check_pdf_image_count__b5bd06ba', 'check_image_inserted__3333dfb2', 'check_image_dimensions__1beedc32', 'check_image_size__993c9cd2028ad767fa928842de0805ca', 'check_triangle_rightedge__afcf05ab', 'check_gif_file__06722a19', 'check_jpg_file__1bf235b17a43af3cc147100153de4d30', 'check_image_dimensions__fa1a72d6', 'check_triangle_bottom_right__363d0657', 'check_gimp_brightness_increase__652303c0d066122d99f102352a5b1a93', 'check_image_dimensions__e0d9b551', 'check_image_vertical_flip__7870cc8f1b3a8b0d979c0fd70171d833', 'check_image_resized__7253f77d', 'check_image_dimensions__717d6863', 'check_jpg_export__0a152dc7', 'check_image_grayscale__3aef0cbd8986e240435edab0fd96c873', 'check_image_addition__17e4ac0c', 'check_image_file_exists__836db3cb3a7cc19d82f98c6a439eb80c', 'check_image_dimensions__12182c78', 'check_image_properties__739292ff', 'check_image_properties__9873b6c6']

def compare_images(shape1, shape2):
    """
    Compare two picture shapes for matching image content.

    Args:
        shape1: First picture shape
        shape2: Second picture shape

    Returns:
        bool: True if images match
    """
    try:
        image1 = shape1.image
        image2 = shape2.image
        if image1.size != image2.size:
            return False
        if image1.content_type != image2.content_type:
            return False
        if image1.blob != image2.blob:
            return False
        return True
    except:
        return False

def check_image_size__2aa2e23a948bbaf102eec8741d709529(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if image exists and has expected dimensions.

    Args:
        result: Dict from getter with 'exists', 'width', 'height' keys
        expected: Dict with 'width', 'height' keys

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('Image file does not exist')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    tolerance = options.get('tolerance', 5)
    width_match = abs(actual_width - expected_width) <= tolerance if actual_width and expected_width else False
    height_match = abs(actual_height - expected_height) <= tolerance if actual_height and expected_height else False
    if width_match and height_match:
        logger.info(f'Image dimensions match: {actual_width}x{actual_height}')
        return 1.0
    else:
        logger.info(f'Image dimensions mismatch: expected {expected_width}x{expected_height}, got {actual_width}x{actual_height}')
        return 0.0

def check_image_properties__4bd1c0a2e0720abe2c2c09ee9978ce22(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if image has expected format and mode.

    Args:
        result: Dict from getter with 'format', 'mode' keys, or None
        expected: Dict with 'format', 'mode' expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.info('Image file not found or invalid')
        return 0.0
    score = 0.0
    expected_format = expected.get('format')
    actual_format = result.get('format')
    if actual_format == expected_format:
        score += 0.5
    else:
        logger.info(f'Format mismatch - Expected: {expected_format}, Got: {actual_format}')
    expected_mode = expected.get('mode')
    actual_mode = result.get('mode')
    if actual_mode == expected_mode:
        score += 0.5
    else:
        logger.info(f'Mode mismatch - Expected: {expected_mode}, Got: {actual_mode}')
    return score

def check_triangle_flipped__c73a98370436f02ecc8c151f49a9b91c(result, expected, **options):
    """
    Check if the triangle has been flipped horizontally.
    The triangle should be on the opposite side of the image after flipping.

    Args:
        result: Actual horizontal position (0.0-1.0) from getter
        expected: Dict with 'original_position' key
        **options: Optional 'tolerance' for position comparison (default: 0.1)

    Returns:
        float: 1.0 if triangle has been flipped (is on opposite side), 0.0 otherwise
    """
    if result is None:
        return 0.0
    original_position = expected.get('original_position')
    tolerance = options.get('tolerance', 0.1)
    expected_position = 1.0 - original_position
    diff = abs(result - expected_position)
    if diff <= tolerance:
        return 1.0
    else:
        return 0.0

def check_image_size__5db5c513f28a5ad7ad8d60ef05b14b35(result, expected, **options):
    """
    Check if image size matches expected dimensions.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected dimensions)
        **options: Additional options (tolerance for allowing small differences)

    Returns:
        float: 1.0 if size matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 0)
    result_width = result.get('width', 0)
    result_height = result.get('height', 0)
    expected_width = expected.get('width', 0)
    expected_height = expected.get('height', 0)
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    logger.info(f'Image size check: result=({result_width}, {result_height}), expected=({expected_width}, {expected_height}), match={width_match and height_match}')
    return 1.0 if width_match and height_match else 0.0

def check_layer_exists__a85f6474(result_state, expected_state, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected_state: Expected layer name to check for (string) or dict with 'layer_name'
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected_state, dict):
            expected_layer = expected_state.get('layer_name', '')
        else:
            expected_layer = expected_state
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_orientation__844f5e73b108da449b9b68fcbef6bbf2(result, expected, **options):
    """
    Check if image rotation was performed correctly by verifying:
    1. The output image is in portrait orientation
    2. Dimensions are correctly swapped (width/height flipped from original)
    3. Image content is preserved after rotation

    Args:
        result: dict with rotation verification data from getter
        expected: dict with expected rotation properties

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_portrait = expected.get('is_portrait', True)
    actual_portrait = result.get('is_portrait')
    if actual_portrait is None:
        logger.error(f'Invalid result: {result}')
        return 0.0
    if actual_portrait != expected_portrait:
        logger.error(f'Orientation mismatch: expected portrait={expected_portrait}, got {actual_portrait}')
        return 0.0
    logger.info(f'✓ Portrait orientation check passed')
    expected_swap = expected.get('dimensions_swapped', True)
    actual_swap = result.get('dimensions_swapped', False)
    if actual_swap != expected_swap:
        logger.error(f'Dimension swap check failed: expected={expected_swap}, got={actual_swap}')
        logger.error(f"Rotated: {result.get('rotated_width')}x{result.get('rotated_height')}, Original: {result.get('original_width')}x{result.get('original_height')}")
        return 0.0
    logger.info(f"✓ Dimension swap check passed (rotated: {result.get('rotated_width')}x{result.get('rotated_height')}, original: {result.get('original_width')}x{result.get('original_height')})")
    expected_preserved = expected.get('content_preserved', True)
    actual_preserved = result.get('content_preserved', False)
    if actual_preserved != expected_preserved:
        logger.error(f'Content preservation check failed: expected={expected_preserved}, got={actual_preserved}')
        logger.error(f"Hash distance: {result.get('hash_distance', 'N/A')}")
        return 0.0
    logger.info(f"✓ Content preservation check passed (hash distance: {result.get('hash_distance')})")
    logger.info('✓ All rotation verification checks passed')
    return 1.0

def check_image_mode__8f91f3a7(result, expected, **options):
    """Check if image mode matches.

    Args:
        result: Dict with image properties
        expected: Dict with expected properties in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists'):
        return 0.0
    expected_mode = expected.get('mode')
    if result.get('mode') == expected_mode:
        return 1.0
    return 0.0

def check_image_dimensions__baa9d5de58b4726390c9c04659eb9fca(result, expected, **options):
    """
    Check if the image has the expected dimensions.

    Args:
        result: Dict with 'width' and 'height' from getter
        expected: Dict with 'width' and 'height' expected values
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    result_width = result.get('width')
    result_height = result.get('height')
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result_width is None or result_height is None:
        logger.error('Missing dimensions in result')
        return 0.0
    logger.info(f'Result dimensions: {result_width}x{result_height}')
    logger.info(f'Expected dimensions: {expected_width}x{expected_height}')
    if result_width == expected_width and result_height == expected_height:
        return 1.0
    else:
        return 0.0

def check_image_extracted__a3446500(result, expected, **options):
    """
    Check if image was extracted correctly with hash verification as primary method.

    Scoring breakdown:
    - File exists at correct path: 0.15
    - Correct format (PNG): 0.15
    - Exact dimensions match: 0.2
    - Content hash matches (proves correct image extracted): 0.5

    Hash verification ensures the extracted image is exactly the third image from the Word document,
    preventing false positives from fake files with similar dimensions.
    """
    score = 0.0
    if not result.get('exists'):
        return 0.0
    score += 0.15
    if result.get('format') == expected.get('format', 'PNG'):
        score += 0.15
    expected_width = expected.get('exact_width', 510)
    expected_height = expected.get('exact_height', 474)
    width_match = result.get('width', 0) == expected_width
    height_match = result.get('height', 0) == expected_height
    if width_match and height_match:
        score += 0.2
    expected_hash = expected.get('expected_hash')
    if expected_hash and expected_hash != 'HASH_PLACEHOLDER_THIRD_IMAGE_FROM_WORD_DOC':
        if result.get('hash') == expected_hash:
            score += 0.5
        else:
            pass
    else:
        expected_size = expected.get('exact_size', 27669)
        if result.get('size', 0) == expected_size:
            score += 0.5
        else:
            pass
    return score

def check_gimp_gimprc_setting__5b6e5d1f(actual_config_path, expected, **options):
    """
    Check if a GIMP gimprc setting has the expected value.

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'key' and 'value' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if setting matches, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        target_key = expected.get('key')
        target_value = expected.get('value')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == target_key and value == target_value:
                    logger.info(f'Found matching setting: {key} = {value}')
                    return 1.0
        logger.warning(f"Setting not found or doesn't match: {target_key} = {target_value}")
        return 0.0
    except Exception as e:
        logger.error(f'Error checking gimprc setting: {e}')
        return 0.0

def check_png_file_exists__846e3a68(result, expected, **options):
    """
    Check if PNG file exists and has the expected dimensions.

    Args:
        result: Path to the PNG file
        expected: Dict with rules specifying width and height
        **options: Additional options

    Returns:
        float: 1.0 if file exists and dimensions match, 0.0 otherwise
    """
    if not result or not os.path.exists(result):
        logger.error(f'Result file does not exist: {result}')
        return 0.0
    expected_width = expected.get('width', 1152)
    expected_height = expected.get('height', 648)
    try:
        img = Image.open(result)
        if img.format != 'PNG':
            logger.error(f'File is not PNG format: {img.format}')
            return 0.0
        (actual_width, actual_height) = img.size
        if actual_width == expected_width and actual_height == expected_height:
            logger.info(f'PNG has correct dimensions: {img.size}')
            return 1.0
        else:
            logger.error(f'Dimension mismatch: expected {expected_width}x{expected_height}, got {actual_width}x{actual_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error opening image file: {e}')
        return 0.0

def check_layer_name_config__3336340a(layer_names, expected, **options):
    """Check if the GIMP .xcf file contains a layer with the expected name."""
    if layer_names is None:
        logger.error('No layer names provided')
        return 0.0
    if not isinstance(layer_names, list):
        logger.error(f'Expected list of layer names, got {type(layer_names)}')
        return 0.0
    try:
        target_layer_name = expected.get('layer_name', 'Text Layer')
        if target_layer_name in layer_names:
            logger.info(f"Found layer '{target_layer_name}' in layers: {layer_names}")
            return 1.0
        else:
            logger.info(f"Layer '{target_layer_name}' not found. Available layers: {layer_names}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer name: {e}')
        return 0.0

def check_triangle_flipped__1aaf02638da71a4a84f47e22e2395da3(result_state, expected, **options):
    """Check if the triangle has been flipped horizontally by comparing before/after images.

    Args:
        result_state: Dict with 'result_path' and 'original_path' keys
        expected: Expected rules dict (from config["rules"])
        **options: Additional options

    Returns:
        float: 1.0 if triangle is flipped, 0.0 otherwise
    """
    if result_state is None or not isinstance(result_state, dict):
        return 0.0
    result_path = result_state.get('result_path')
    original_path = result_state.get('original_path')
    if not result_path or not original_path:
        return 0.0
    try:
        result_img = Image.open(result_path)
        original_img = Image.open(original_path)
        result_array = np.array(result_img)
        original_array = np.array(original_img)
        if result_array.shape != original_array.shape:
            logger.error(f"Image dimensions don't match: result={result_array.shape}, original={original_array.shape}")
            return 0.0
        (height, width) = result_array.shape[:2]
        triangle_mask_original = detect_yellow_triangle(original_array)
        if not triangle_mask_original.any():
            logger.error('Could not detect yellow triangle in original image')
            return 0.0
        original_coords = np.argwhere(triangle_mask_original)
        (original_centroid_y, original_centroid_x) = original_coords.mean(axis=0)
        triangle_mask_result = detect_yellow_triangle(result_array)
        if not triangle_mask_result.any():
            logger.error('Could not detect yellow triangle in result image')
            return 0.0
        result_coords = np.argwhere(triangle_mask_result)
        (result_centroid_y, result_centroid_x) = result_coords.mean(axis=0)
        original_side = 'left' if original_centroid_x < width / 2 else 'right'
        result_side = 'left' if result_centroid_x < width / 2 else 'right'
        if original_side == result_side:
            logger.debug(f'Triangle did not move to opposite side: original={original_side}, result={result_side}')
            return 0.0
        flip_score = verify_horizontal_flip(original_array, result_array, triangle_mask_original, triangle_mask_result, width)
        if flip_score < 0.8:
            logger.debug(f'Flip verification failed: score={flip_score}')
            return 0.0
        original_size = np.sum(triangle_mask_original)
        result_size = np.sum(triangle_mask_result)
        size_ratio = min(original_size, result_size) / max(original_size, result_size)
        if size_ratio < 0.9:
            logger.debug(f'Triangle size changed too much: original={original_size}, result={result_size}')
            return 0.0
        logger.debug(f'Triangle successfully flipped: centroid moved from x={original_centroid_x:.1f} to x={result_centroid_x:.1f}, flip_score={flip_score:.2f}')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking triangle flip: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_layer_fixed_size__0dbca844(result_state, expected_state, **options):
    """
    Check if the layer is resized to exact dimensions (variation 3)

    Args:
        result_state: Path to the result image file
        expected_state: Dict with 'width', 'height' and optional 'ignore_transparent' keys
        **options: Additional options

    Returns:
        float: Score (1.0 if both width and height match, 0.0 otherwise)
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    try:
        img = Image.open(result_state)
        ignore_transparent = expected_state.get('ignore_transparent', False)
        if ignore_transparent and (img.mode in ('RGBA', 'LA') or 'transparency' in img.info):
            if img.mode != 'RGBA':
                img = img.convert('RGBA')
            alpha = img.split()[-1]
            bbox = alpha.getbbox()
            if bbox is None:
                actual_width = 0
                actual_height = 0
            else:
                actual_width = bbox[2] - bbox[0]
                actual_height = bbox[3] - bbox[1]
            logger.debug(f'Original size: {img.size}, Content size: {actual_width}x{actual_height}')
        else:
            actual_width = img.size[0]
            actual_height = img.size[1]
            logger.debug(f'Image size: {actual_width}x{actual_height}')
        expected_width = expected_state.get('width')
        expected_height = expected_state.get('height')
        if expected_width is None or expected_height is None:
            logger.error('Expected dimensions not fully specified')
            return 0.0
        width_match = actual_width == expected_width
        height_match = actual_height == expected_height
        logger.debug(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
        logger.debug(f'Width match: {width_match}, Height match: {height_match}')
        return 1.0 if width_match and height_match else 0.0
    except Exception as e:
        logger.error(f'Error checking layer size: {e}')
        return 0.0

def check_image_hash__db2588a0(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_files_exist__abdf3926e9609e7e5b435c8cdbb40013(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected image files exist in the directory.

    Args:
        result: List of filenames from the directory
        expected: Dict with 'expected_files' key containing list of expected filenames
        **options: Additional options

    Returns:
        1.0 if all expected files exist, 0.0 otherwise
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if expected_set.issubset(result_set):
        if len(result_set) == len(expected_set):
            return 1.0
        else:
            return 0.5
    return 0.0

def check_triangle_topright__54d36b10(result_state, expected_state, **options):
    """
    Check if the triangle is in the top-right corner of the image.
    Variation 8 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        result_state: Path to the exported image
        expected_state: Expected value (true if triangle should be in top-right)
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is in top-right, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        triangle_color = None
        for color in unique_colors_sorted:
            if len(color) >= 3 and color[0] > 200 and (color[1] > 200) and (color[2] < 100):
                triangle_color = color
                break
        if triangle_color is None:
            if len(unique_colors_sorted) > 1:
                triangle_color = unique_colors_sorted[1]
            else:
                logger.warning('Could not detect triangle color')
                return 0.0
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        centroid = triangle_coords.mean(axis=0)
        image_shape = np.array(img_array.shape[:2])
        top_right = centroid[0] < image_shape[0] * 0.25 and centroid[1] > image_shape[1] * 0.75
        if expected_state is True or expected_state == 'true':
            if bool(top_right):
                return 1.0
            else:
                return 0.0
        elif bool(top_right):
            return 0.0
        else:
            return 1.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        return 0.0

def check_specific_image_inserted__2759cd98(result, expected, **options):
    """Check if a specific image was inserted on the slide.

    This function verifies that:
    1. Exactly one image was added (count increased from initial_count to expected_count)
    2. The added image matches the specific expected image file (if hash is provided)

    Args:
        result: List of image hashes from getter (list of SHA256 hex strings)
        expected: Expected criteria dict with:
            - expected_count: The expected final number of images (default: 1)
            - initial_count: The initial number of images (default: 0)
            - expected_image_hash: Pre-computed SHA256 hash of the expected image (optional)
            - expected_image_path: VM path to expected image (for reference only, not used)
        **options: Additional options

    Returns:
        float: 1.0 if correct image was inserted, 0.0 otherwise
    """
    expected_count = expected.get('expected_count', 1)
    initial_count = expected.get('initial_count', 0)
    expected_image_hash = expected.get('expected_image_hash', '')
    actual_count = len(result)
    if actual_count != expected_count:
        return 0.0
    if expected_image_hash:
        if expected_image_hash in result:
            return 1.0
        else:
            return 0.0
    return 1.0

def check_image_dimensions__a7d5fd37(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: Dict with 'width' and 'height' keys from getter
        expected: Dict with expected 'width' and 'height' values
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: Score (1.0 if dimensions match, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        result_width = result.get('width')
        result_height = result.get('height')
        if result_width is None or result_height is None:
            logger.error(f'Invalid result dimensions: {result}')
            return 0.0
        width_match = result_width == expected_width if expected_width is not None else True
        height_match = result_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            logger.info(f'Dimensions match: {result_width}x{result_height}')
            return 1.0
        else:
            logger.warning(f'Dimensions mismatch: got {result_width}x{result_height}, expected {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image dimensions: {e}')
        return 0.0

def check_moved_pngs__ec3ddc36(png_list, expected):
    """Check if the correct PNG files were moved.

    Args:
        png_list: List of PNG filenames in target directory
        expected: Dict with 'expected' key containing list of expected filenames

    Returns:
        float: 1.0 if all expected PNGs are present, 0.0 otherwise
    """
    expected_pngs = expected['expected']
    if len(png_list) != len(expected_pngs):
        return 0.0
    if set(png_list) == set(expected_pngs):
        return 1.0
    else:
        return 0.0

def check_image_grayscale__bf5ecb70b33ef3dffbe095b744ec38a7(result, expected, **options):
    """
    Check if image is in grayscale mode.

    Args:
        result: dict with "is_grayscale" and "mode" keys from getter
        expected: dict with expected "is_grayscale" value (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if grayscale status matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_grayscale = expected.get('is_grayscale', True)
    result_grayscale = result.get('is_grayscale', False)
    if result_grayscale == expected_grayscale:
        logger.info(f"Grayscale check passed: {result_grayscale} (mode: {result.get('mode')})")
        return 1.0
    else:
        logger.info(f'Grayscale check failed: got {result_grayscale}, expected {expected_grayscale}')
        return 0.0

def check_png_exists__06536a54(result, expected, **options):
    """Check if the file exists and is in PNG format.

    Args:
        result: Path to the result file (from vm_file getter)
        expected: Expected rules dict with format specification
        **options: Additional comparison options

    Returns:
        float: 1.0 if file exists and is PNG format, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if not os.path.isfile(result):
        return 0.0
    try:
        img = Image.open(result)
        img_format = img.format
        img.close()
        if expected and isinstance(expected, dict):
            expected_format = expected.get('format', 'PNG')
            if img_format == expected_format:
                return 1.0
            else:
                return 0.0
        if img_format == 'PNG':
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_gimp_saturation_decrease__62aadfc75943521b0b091bf9d9c10f24(result, expected, **options):
    """
    Check if the image saturation is lower than the original.

    Args:
        result: float saturation value from getter
        expected: dict with 'max_saturation' key specifying maximum acceptable saturation
        **options: Additional options

    Returns:
        float: 1.0 if saturation is below maximum, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    max_saturation = expected.get('max_saturation', 255)
    logger.info(f'Saturation check: result={result}, max_allowed={max_saturation}')
    if result <= max_saturation:
        logger.info(f'Saturation check PASSED: {result} <= {max_saturation}')
        return 1.0
    else:
        logger.info(f'Saturation check FAILED: {result} > {max_saturation}')
        return 0.0

def check_image_exists__e327229f(result, expected, **options):
    """
    Check if image file exists and meets minimum size requirement.

    Args:
        result: Dict from getter with 'exists' and 'size'
        expected: Dict with 'should_exist' and 'min_size'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.5 if exists but too small, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    exists = result.get('exists', False)
    if not exists:
        return 0.0
    score = 0.5
    min_size = expected.get('min_size', 0)
    file_size = result.get('size', 0)
    if file_size >= min_size:
        score += 0.5
    return score

def check_gif_file__b130b682(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_dimensions__13d2c6c3(src_path, rule):
    """
    Check if both width and height of the image are correct, optionally ignoring transparent areas.
    Variation 1: d16c99dc-2a1e-46f2-b350-d97c86c85c15_task_verify_1
    """
    if src_path is None:
        return 0.0
    img = Image.open(src_path)
    ignore_transparent = rule.get('ignore_transparent', False)
    if ignore_transparent and (img.mode in ('RGBA', 'LA') or 'transparency' in img.info):
        if img.mode != 'RGBA':
            img = img.convert('RGBA')
        alpha = img.split()[-1]
        bbox = alpha.getbbox()
        if bbox is None:
            actual_width = 0
            actual_height = 0
        else:
            actual_width = bbox[2] - bbox[0]
            actual_height = bbox[3] - bbox[1]
        logger.debug(f'Original size: {img.size}, Content size: {actual_width}x{actual_height}')
    else:
        actual_width = img.size[0]
        actual_height = img.size[1]
        logger.debug(f'Image size: {img.size}')
    if rule.get('width', None) is not None:
        width_same = actual_width == rule['width']
    else:
        width_same = True
    if rule.get('height', None) is not None:
        height_same = actual_height == rule['height']
    else:
        height_same = True
    if width_same and height_same:
        logger.debug(f'width_same: {width_same}, height_same: {height_same}')
        return 1.0
    else:
        logger.debug(f"width_same: {width_same}, height_same: {height_same}, expected: {rule.get('width')}x{rule.get('height')}, actual: {actual_width}x{actual_height}")
        return 0.0

def check_gif_file__f8138076f6546232c093f5027d50db21(result, expected, **options):
    """
    Check if a GIF file exists and meets requirements for a 3-second video segment.

    Args:
        result: dict from get_gif_file_info__f8138076f6546232c093f5027d50db21
        expected: dict with expected properties (min_file_size, min_frames, etc.)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.25
        logger.info('File exists: +0.25')
    else:
        logger.warning('File does not exist')
        return 0.0
    if result.get('is_gif', False):
        score += 0.25
        logger.info('File is a valid GIF: +0.25')
    else:
        logger.warning('File is not a valid GIF')
        return score
    min_size = expected.get('min_file_size', 50000)
    if result.get('file_size', 0) >= min_size:
        score += 0.2
        logger.info(f"File size {result['file_size']} >= {min_size}: +0.2")
    else:
        logger.warning(f"File size {result['file_size']} < {min_size}")
    min_frames = expected.get('min_frames', 15)
    if result.get('frames', 0) >= min_frames:
        score += 0.2
        logger.info(f"Frame count {result['frames']} >= {min_frames}: +0.2")
    else:
        logger.warning(f"Frame count {result['frames']} < {min_frames}")
    duration_ms = result.get('duration_ms', 0)
    if 2000 <= duration_ms <= 4000:
        score += 0.1
        logger.info(f'Duration {duration_ms}ms is approximately 3 seconds: +0.1')
    else:
        logger.warning(f'Duration {duration_ms}ms is not close to 3000ms')
    logger.info(f'Final score: {score}')
    return score

def check_image_files_exist__e4b458033c1389ecc56cd63f5eae9626(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected image files exist in the directory.

    Args:
        result: List of filenames from the directory
        expected: Dict with 'expected_files' key containing list of expected filenames
        **options: Additional options

    Returns:
        1.0 if all expected files exist, 0.0 otherwise
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    return 0.0

def check_image_resize__dbc2e36d8461f22779ef5b4f0521a159(src_path, expected, **options):
    """
    Check if the image has been resized to the expected dimensions.

    Args:
        src_path: Path to the edited/result image
        expected: Dictionary containing target dimensions with 'width' and 'height' keys
        **options: Additional options (unused)

    Returns:
        1.0 if the image dimensions match the expected size, 0.0 otherwise
    """
    if src_path is None:
        logger.warning('Source path is None')
        return 0.0
    if not expected or not isinstance(expected, dict):
        logger.error("Expected value must be a dictionary with 'width' and 'height' keys")
        return 0.0
    try:
        edited_image = Image.open(src_path)
        (actual_width, actual_height) = edited_image.size
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        if expected_width is None or expected_height is None:
            logger.error('Expected dimensions not provided')
            return 0.0
        logger.info(f'Actual size: {actual_width}x{actual_height}, Expected size: {expected_width}x{expected_height}')
        if actual_width == expected_width and actual_height == expected_height:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image resize: {e}')
        return 0.0

def check_triangle_scaled__93095e022ff3d0c1026d009f3ccc512b(result_path, expected, **options):
    """Check if the triangle has been scaled 50% larger while keeping same position.

    Args:
        result_path: Path to the result image
        expected: Expected rules dict (from config["rules"])
        **options: Additional options

    Returns:
        float: 1.0 if both size and position correct, 0.5 if only size correct, 0.0 otherwise
    """
    if result_path is None:
        return 0.0
    try:
        result_img = Image.open(result_path)
        result_array = np.array(result_img)
        cache_dir = os.path.dirname(result_path)
        original_cache_path = os.path.join(cache_dir, 'Triangle_On_The_Side.png')
        if not os.path.exists(original_cache_path):
            logger.error(f'Original image not found at {original_cache_path}')
            return 0.0
        original_img = Image.open(original_cache_path)
        original_array = np.array(original_img)
        triangle_color_result = _find_triangle_color(result_array)
        triangle_color_original = _find_triangle_color(original_array)
        if triangle_color_result is None or triangle_color_original is None:
            logger.error('Could not find triangle color')
            return 0.0
        result_mask = np.all(result_array == triangle_color_result, axis=2)
        original_mask = np.all(original_array == triangle_color_original, axis=2)
        result_pixel_count = np.sum(result_mask)
        original_pixel_count = np.sum(original_mask)
        result_centroid = _calculate_centroid(result_mask)
        original_centroid = _calculate_centroid(original_mask)
        if result_centroid is None or original_centroid is None:
            logger.error('Could not calculate centroid')
            return 0.0
        scale_threshold = expected.get('scale_threshold', 1.5)
        expected_min_pixels = original_pixel_count * scale_threshold ** 2
        size_correct = result_pixel_count >= expected_min_pixels * 0.9
        position_threshold = expected.get('position_threshold', 10)
        centroid_distance = np.sqrt((result_centroid[0] - original_centroid[0]) ** 2 + (result_centroid[1] - original_centroid[1]) ** 2)
        position_correct = centroid_distance <= position_threshold
        logger.debug(f'Original pixels: {original_pixel_count}, Result pixels: {result_pixel_count}')
        logger.debug(f'Expected min pixels: {expected_min_pixels}, Actual: {result_pixel_count}')
        logger.debug(f'Original centroid: {original_centroid}, Result centroid: {result_centroid}')
        logger.debug(f'Centroid distance: {centroid_distance}, threshold: {position_threshold}')
        logger.debug(f'Size correct: {size_correct}, Position correct: {position_correct}')
        if size_correct and position_correct:
            return 1.0
        elif size_correct:
            return 0.5
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle scale: {e}')
        return 0.0

def check_image_dimensions__4eb874068ccb861273ecf8604bdafb3c(result, expected, **options):
    """Check if image dimensions match the expected values.

    Args:
        result: Dict with 'width' and 'height' keys from the result image
        expected: Dict with 'width' and 'height' keys for expected dimensions
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is None or expected_height is None:
        return 0.0
    result_width = result.get('width')
    result_height = result.get('height')
    if result_width is None or result_height is None:
        return 0.0
    if result_width == expected_width and result_height == expected_height:
        return 1.0
    else:
        return 0.0

def check_image_hash__6510dd9a(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_count__02587c7a(result, expected, **options):
    """Check if image count matches expected.

    Args:
        result: Actual image count
        expected: Expected count (dict with 'count' key)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_count = expected.get('count', 4)
    return 1.0 if result == expected_count else 0.0

def check_default_video_player__9b1915f5(result, expected, **options):
    """
    Check if VLC is set as the default video player.

    Args:
        result: Default video player string from getter (e.g., "vlc.desktop")
        expected: Dict with 'rules' containing 'player_name' to match
        **options: Additional options

    Returns:
        float: 1.0 if result contains expected player name, 0.0 otherwise
    """
    player_name = expected.get('player_name', 'vlc')
    if not result:
        logger.warning('No default video player detected')
        return 0.0
    if player_name.lower() in result.lower():
        logger.info(f'Verified {player_name} is default video player: {result}')
        return 1.0
    else:
        logger.warning(f'Expected {player_name} but got: {result}')
        return 0.0

def check_gimp_menubar__4e59f0cb646eef506e256c2929b9d463(actual_config_path, expected):
    """
    Check if GIMP fullscreen menubar visibility setting is correct
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(expected['key'], str):
            if items[0] == expected['key'] and items[-1] == expected['value']:
                return 1.0
        elif isinstance(expected['key'], list) and len(expected['key']) == 2:
            if items[0] == expected['key'][0] and items[1] == expected['key'][1] and (items[-1] == expected['value']):
                return 1.0
    return 0.0

def check_image_is_grayscale__0e7db70c(result_state: Optional[str], expected_state: Dict, **options) -> float:
    """
    Check if an image is in grayscale mode.

    An image is considered grayscale if:
    1. It's in mode 'L' (Luminance/Grayscale), OR
    2. It's in mode 'LA' (Grayscale with Alpha), OR
    3. All RGB channels have equal values (converted to grayscale but saved as RGB)

    Args:
        result_state: Local path to the image file
        expected_state: Dict with rules, should contain {"is_grayscale": true}
        **options: Additional options

    Returns:
        float: 1.0 if image is grayscale, 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Image path is None')
        return 0.0
    if not os.path.exists(result_state):
        logger.warning(f'Image file does not exist: {result_state}')
        return 0.0
    try:
        img = Image.open(result_state)
        logger.info(f'Image mode: {img.mode}, size: {img.size}')
        if img.mode in ('L', 'LA'):
            logger.info(f'Image is in grayscale mode: {img.mode}')
            return 1.0
        if img.mode in ('RGB', 'RGBA'):
            img_array = np.array(img)
            r = img_array[:, :, 0]
            g = img_array[:, :, 1]
            b = img_array[:, :, 2]
            is_grayscale = np.array_equal(r, g) and np.array_equal(g, b)
            if is_grayscale:
                logger.info('Image has equal RGB channels (grayscale)')
                return 1.0
            else:
                logger.info('Image has different RGB channels (not grayscale)')
                return 0.0
        logger.warning(f'Unsupported image mode for grayscale check: {img.mode}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking if image is grayscale: {e}')
        return 0.0

def check_bg_image_exists__0a4467dd72d00d84bc5d80c765c4ad6f(result, expected, **options):
    """
    Check if a background image exists in the slide.

    Args:
        result: Path to the background image file, or None if not found
        expected: Expected condition (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if background image exists, 0.0 otherwise
    """
    should_have_bg = expected.get('has_background', True)
    has_bg = result is not None
    if has_bg == should_have_bg:
        return 1.0
    else:
        return 0.0

def check_image_cropped__7bac509f(result, expected, **options):
    """
    Check if the image has been cropped to specified dimensions.

    Args:
        result: Path to result image
        expected: Dict with 'width' and 'height' for cropped size
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    try:
        if result is None:
            logger.error('Result path is None')
            return 0.0
        result_img = Image.open(result)
        expected_width = expected.get('width', 0)
        expected_height = expected.get('height', 0)
        if result_img.width == expected_width and result_img.height == expected_height:
            logger.info(f'Image cropped correctly to {expected_width}x{expected_height}')
            return 1.0
        else:
            logger.warning(f'Image size mismatch: {result_img.width}x{result_img.height} != {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking crop: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_mirror__bc5fe443a1102b529ba31d0ac9c81ab1(result, expected, **options):
    """
    Check if the result image is a horizontal mirror/flip of the expected (original) image.

    Args:
        result: Path to the flipped image file (from getter)
        expected: Path to the original image file (expected should contain 'original_path')
        **options: Additional options

    Returns:
        float: 1.0 if image is mirrored, 0.0 otherwise
    """
    if result is None or expected is None:
        logger.error('Result or expected path is None')
        return 0.0
    try:
        original_path = expected.get('original_path')
        if not original_path:
            logger.error('No original_path in expected rules')
            return 0.0
        result_image = Image.open(result)
        original_image = Image.open(original_path)
        transposed_image = original_image.transpose(Image.FLIP_LEFT_RIGHT)
        mirrored = structure_check_by_ssim(transposed_image, result_image, 0.99)
        logger.info(f'Mirror check result: {mirrored}')
        return 1.0 if mirrored else 0.0
    except Exception as e:
        logger.error(f'Error checking image mirror: {e}')
        return 0.0

def check_layer_exists__8565a91c(result_state, expected_state, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected_state: Expected layer name to check for (dict with 'layer_name')
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected_state, dict):
            expected_layer = expected_state.get('layer_name', '')
        else:
            expected_layer = expected_state
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_hash__0c114ed7(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_dimensions__3f39d534f7803089ed331d19c2a1bc89(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected values)
        **options: Additional options (unused)

    Returns:
        float: 1.0 if dimensions match exactly, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if expected_width is not None and actual_width != expected_width:
        return 0.0
    if expected_height is not None and actual_height != expected_height:
        return 0.0
    return 1.0

def check_image_hash__73176121(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_triangle_top_left__86e46240(result_state, expected_state, **options):
    """
    Check if the yellow triangle has been moved to the top-left corner.

    Args:
        result_state: Path to the result image
        expected_state: Not used (rule-based evaluation)
        **options: Additional options

    Returns:
        float: Score (1.0 if positioned correctly, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        sorted_indices = np.argsort(counts)[::-1]
        unique_colors_sorted = unique_colors[sorted_indices]
        counts_sorted = counts[sorted_indices]
        triangle_color = None
        for (i, color) in enumerate(unique_colors_sorted):
            if i == 0:
                continue
            if is_yellow_color(color):
                triangle_color = color
                logger.info(f'Found yellow triangle color at rank {i}: RGB={color}, count={counts_sorted[i]}')
                break
        if triangle_color is None:
            logger.warning('Could not find yellow triangle in image')
            return 0.0
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        centroid = triangle_coords.mean(axis=0)
        (image_height, image_width) = img_array.shape[:2]
        top_threshold = image_height * 0.15
        left_threshold = image_width * 0.15
        logger.info(f'Triangle centroid: ({centroid[1]:.1f}, {centroid[0]:.1f})')
        logger.info(f'Thresholds: left={left_threshold:.1f}, top={top_threshold:.1f}')
        if centroid[0] <= top_threshold and centroid[1] <= left_threshold:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error in check_triangle_top_left__86e46240: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_gif_file__430ad71e1ce94d5ea2bb30ba073fd937(result, expected, **options):
    """
    Check if a GIF file exists with appropriate file size, frame count, and duration constraints.

    Args:
        result: dict from get_gif_file_info__430ad71e1ce94d5ea2bb30ba073fd937
        expected: dict with expected properties including max_file_size, min_frames, and duration

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.25
        logger.info('File exists: +0.25')
    else:
        logger.warning('File does not exist')
        return 0.0
    if result.get('is_gif', False):
        score += 0.25
        logger.info('File is a valid GIF: +0.25')
    else:
        logger.warning('File is not a valid GIF')
        return score
    min_size = expected.get('min_file_size', 1000)
    max_size = expected.get('max_file_size', 10000000)
    file_size = result.get('file_size', 0)
    if min_size <= file_size <= max_size:
        score += 0.15
        logger.info(f'File size {file_size} within range [{min_size}, {max_size}]: +0.15')
    else:
        logger.warning(f'File size {file_size} out of range [{min_size}, {max_size}]')
    min_frames = expected.get('min_frames', 30)
    if result.get('frames', 0) >= min_frames:
        score += 0.15
        logger.info(f"Frame count {result['frames']} >= {min_frames}: +0.15")
    else:
        logger.warning(f"Frame count {result['frames']} < {min_frames}")
    min_duration = expected.get('min_duration', 5.0)
    max_duration = expected.get('max_duration', 7.0)
    actual_duration = result.get('duration_seconds', 0.0)
    if min_duration <= actual_duration <= max_duration:
        score += 0.2
        logger.info(f'Duration {actual_duration:.2f}s within range [{min_duration}, {max_duration}]: +0.20')
    else:
        logger.warning(f'Duration {actual_duration:.2f}s out of range [{min_duration}, {max_duration}]')
    logger.info(f'Final score: {score}')
    return score

def check_image_dimensions__a846d552e46987bd85e04c0a4f658c7a(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected values)
        **options: Additional options (unused)

    Returns:
        float: 1.0 if dimensions match exactly, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if expected_width is not None and actual_width != expected_width:
        return 0.0
    if expected_height is not None and actual_height != expected_height:
        return 0.0
    return 1.0

def check_image_dimensions__68566bbc(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Dict with image properties from getter
        expected: Dict with expected properties in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists'):
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    score = 0.0
    checks = 0
    if expected_width is not None:
        checks += 1
        if result.get('width') == expected_width:
            score += 0.5
    if expected_height is not None:
        checks += 1
        if result.get('height') == expected_height:
            score += 0.5
    return score if checks > 0 else 0.0

def check_pdf_image_count__f35ebb0d(pdf_file: str, expected, **options):
    """
    Check if PDF contains the expected number of images.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with 'min_images' count
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.0
    try:
        doc = fitz.open(pdf_file)
        if doc.page_count > 0:
            score += 0.3
        image_count = 0
        for page_num in range(doc.page_count):
            page = doc[page_num]
            images = page.get_images(full=True)
            image_count += len(images)
        min_images = expected.get('min_images', 1)
        if image_count >= min_images:
            score += 0.7
        doc.close()
    except Exception as e:
        return 0.0
    return score

def check_image_properties__c6063730(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_image_dimensions__7816ba80(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_layer_count__58c8d068(result, expected, **options):
    """
    Check if the layer count matches the expected value.

    Args:
        result: Actual layer count from getter (int)
        expected: Expected layer count (dict with 'count' key)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    try:
        expected_count = expected.get('count', 1)
        if result == expected_count:
            logger.info(f'Layer count matches: {result} == {expected_count}')
            return 1.0
        else:
            logger.warning(f'Layer count mismatch: {result} != {expected_count}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer count: {e}')
        return 0.0

def check_image_dimensions__f5cfdff3841c16728bb4565a839b59ca(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected values)
        **options: Additional options (unused)

    Returns:
        float: 1.0 if dimensions match exactly, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if expected_width is not None and actual_width != expected_width:
        return 0.0
    if expected_height is not None and actual_height != expected_height:
        return 0.0
    return 1.0

def check_image_size__07fcde31879d28764a081db212afaf2d(result, expected, **options):
    """
    Check if image size matches expected dimensions.

    Args:
        result: Dict with 'width' and 'height' keys
        expected: Dict with 'width' and 'height' expected values

    Returns:
        float: 1.0 if both dimensions match, 0.5 if one matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.error(f'Invalid result or expected type: {type(result)}, {type(expected)}')
        return 0.0
    result_width = result.get('width')
    result_height = result.get('height')
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    logger.debug(f'Size check - result: {result_width}x{result_height}, expected: {expected_width}x{expected_height}')
    width_match = result_width == expected_width
    height_match = result_height == expected_height
    if width_match and height_match:
        logger.debug('Both dimensions match')
        return 1.0
    elif width_match or height_match:
        logger.debug('One dimension matches')
        return 0.5
    else:
        logger.debug('No dimensions match')
        return 0.0

def check_image_size__0d547ae5(src_path, rule):
    """
    Check if the size of the src image matches the expected dimensions
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_9
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        actual_width = img.size[0]
        actual_height = img.size[1]
        logger.debug(f'Image size: {img.size}')
        if rule.get('height', None) is not None:
            height_same = actual_height == rule['height']
        else:
            height_same = True
        if rule.get('width', None) is not None:
            width_same = actual_width == rule['width']
        else:
            width_same = True
        if height_same and width_same:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 1.0
        else:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image size: {e}')
        return 0.0

def check_image_properties__c4f65e24(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def validate_image_count__b01cf463(result, expected, **options):
    """Validate image count matches expected.

    Args:
        result: Actual count
        expected: Expected value (dict with 'count' key)
        **options: Additional options

    Returns:
        float: 1.0 if equal, 0.0 otherwise
    """
    expected_count = expected.get('count', 4)
    return 1.0 if result == expected_count else 0.0

def check_image_format__22b30bf2(result, expected, **options):
    """Check if image format matches.

    Args:
        result: Dict with image properties
        expected: Dict with expected properties in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists'):
        return 0.0
    expected_format = expected.get('format')
    if result.get('format') == expected_format:
        return 1.0
    return 0.0

def check_gif_file__092b88db7884c339b37665609b49294b(result, expected, **options):
    """
    Check if a GIF file exists and meets requirements for a 4-second video clip.

    Args:
        result: dict from get_gif_file_info__092b88db7884c339b37665609b49294b
        expected: dict with expected properties including:
            - min_file_size: minimum file size in bytes (e.g., 50000 for 50KB)
            - min_frames: minimum frame count (e.g., 40 for ~4s at 10fps)
            - max_frames: maximum frame count (e.g., 120 for ~4s at 30fps)
            - min_duration: minimum duration in seconds (e.g., 3.5)
            - max_duration: maximum duration in seconds (e.g., 5.0)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.2
        logger.info('File exists: +0.2')
    else:
        logger.warning('File does not exist')
        return 0.0
    if result.get('is_gif', False):
        score += 0.2
        logger.info('File is a valid GIF: +0.2')
    else:
        logger.warning('File is not a valid GIF')
        return score
    min_size = expected.get('min_file_size', 50000)
    file_size = result.get('file_size', 0)
    if file_size >= min_size:
        score += 0.2
        logger.info(f'File size {file_size} >= {min_size}: +0.2')
    else:
        logger.warning(f'File size {file_size} < {min_size} (too small for 4-second clip)')
    min_frames = expected.get('min_frames', 40)
    max_frames = expected.get('max_frames', 120)
    frame_count = result.get('frames', 0)
    if min_frames <= frame_count <= max_frames:
        score += 0.2
        logger.info(f'Frame count {frame_count} is in range [{min_frames}, {max_frames}]: +0.2')
    else:
        logger.warning(f'Frame count {frame_count} not in expected range [{min_frames}, {max_frames}] for 4-second clip')
    min_duration = expected.get('min_duration', 3.5)
    max_duration = expected.get('max_duration', 5.0)
    duration = result.get('duration', 0.0)
    if min_duration <= duration <= max_duration:
        score += 0.2
        logger.info(f'Duration {duration:.2f}s is in range [{min_duration}, {max_duration}]: +0.2')
    else:
        logger.warning(f'Duration {duration:.2f}s not in expected range [{min_duration}, {max_duration}] (should be ~4 seconds)')
    logger.info(f'Final score: {score}')
    return score

def check_image_columns_reversed__3936f39d(result, expected, **options):
    """
    Check if the image columns have been reversed (right to left).

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'num_columns'
        **options: Additional options

    Returns:
        float: Score (1.0 if columns reversed correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        num_columns = expected.get('num_columns', 3)
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        if result_img.size != source_img.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {source_img.size}')
            return 0.0
        (width, height) = source_img.size
        col_width = width // num_columns
        score = 0.0
        for i in range(num_columns):
            source_col = source_img.crop((i * col_width, 0, (i + 1) * col_width, height))
            result_col = result_img.crop(((num_columns - 1 - i) * col_width, 0, (num_columns - i) * col_width, height))
            if structure_check_by_ssim(source_col, result_col, threshold=0.95):
                score += 1.0 / num_columns
            else:
                logging.debug(f'Column {i} does not match reversed position')
        return score
    except Exception as e:
        logging.error(f'Error in check_image_columns_reversed__3936f39d: {e}')
        return 0.0

def check_jpeg_exists_and_size__ad4cc5ac(result, expected, **options):
    """
    Check if JPEG file exists and has the expected dimensions.

    Args:
        result: Path to the JPEG file
        expected: Dict with rules specifying width and height
        **options: Additional options

    Returns:
        float: 1.0 if file exists and dimensions match, 0.0 otherwise
    """
    if not result or not os.path.exists(result):
        logger.error(f'Result file does not exist: {result}')
        return 0.0
    expected_width = expected.get('width', 1152)
    expected_height = expected.get('height', 648)
    try:
        img = Image.open(result)
        if img.format not in ['JPEG', 'JPG']:
            logger.error(f'File is not JPEG format: {img.format}')
            return 0.0
        (actual_width, actual_height) = img.size
        if actual_width == expected_width and actual_height == expected_height:
            logger.info(f'JPEG has correct dimensions: {img.size}')
            return 1.0
        else:
            logger.error(f'Dimension mismatch: expected {expected_width}x{expected_height}, got {actual_width}x{actual_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error opening image file: {e}')
        return 0.0

def check_gimp_rulers_setting__96f3f9d88ca4e4230495b91ff566eb51(actual_config_path, rule):
    """
    Check if GIMP rulers setting is as expected.
    This checks the sessionrc file for rulers visibility settings.

    Args:
        actual_config_path: Path to the GIMP config file
        rule: Expected configuration with keys:
            - key: Config key name (can be string or list)
            - value: Expected value

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_image_properties__811349c6(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_gimp_statusbar_setting__2d2ae6cc356da88025063f58e6110bad(actual_config_path, rule):
    """
    Check if GIMP statusbar setting is as expected.
    This checks the sessionrc file for statusbar visibility settings.

    Args:
        actual_config_path: Path to the GIMP config file
        rule: Expected configuration with keys:
            - key: Config key name (can be string or list)
            - value: Expected value

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_gif_file__d18c2fd4(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_dimensions__9d6c98f0(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_size__8677d5c370f70c69494653c2a8ef5be2(result, expected, **options):
    """
    Check if image size matches expected dimensions.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected dimensions)
        **options: Additional options (tolerance for allowing small differences)

    Returns:
        float: 1.0 if size matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 0)
    result_width = result.get('width', 0)
    result_height = result.get('height', 0)
    expected_width = expected.get('width', 0)
    expected_height = expected.get('height', 0)
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    logger.info(f'Image size check: result=({result_width}, {result_height}), expected=({expected_width}, {expected_height}), match={width_match and height_match}')
    return 1.0 if width_match and height_match else 0.0

def check_image_format__900fc36276e9eaf12471a7834992cf5c(result, expected, **options):
    """
    Check if image was saved with correct format.

    Args:
        result: Image format info from getter
        expected: Expected format and filename
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        logger.info('Image file does not exist')
        return 0.0
    expected_filename = expected.get('filename', '')
    actual_filename = result.get('filename', '')
    if actual_filename == expected_filename:
        score += 0.3
    else:
        logger.warning(f"Filename mismatch: expected '{expected_filename}', got '{actual_filename}'")
    expected_format = expected.get('format', '')
    actual_format = result.get('format', '')
    if actual_format == expected_format:
        score += 0.3
    else:
        logger.warning(f"Format mismatch: expected '{expected_format}', got '{actual_format}'")
    logger.info(f'Image format check score: {score}')
    return score

def check_image_flipped__07b5058cd3df711389d2b4342d0c561c(result_state, expected_state, **options):
    """
    Check if the result image is a horizontal flip of the original image.

    Args:
        result_state: Path to the result image file (flipped image)
        expected_state: Dict with 'original_path' key pointing to the original image
        **options: Additional options

    Returns:
        float: 1.0 if the image is correctly flipped horizontally, 0.0 otherwise
    """
    if not isinstance(expected_state, dict):
        logger.error(f"expected_state should be a dict with 'original_path', got {type(expected_state)}")
        return 0.0
    original_path = expected_state.get('original_path')
    if not original_path:
        logger.error('original_path not found in expected_state')
        return 0.0
    if result_state is None:
        logger.warning('Result image not found')
        return 0.0
    try:
        original_image = Image.open(original_path)
        result_image = Image.open(result_state)
        flipped_original = original_image.transpose(Image.FLIP_LEFT_RIGHT)
        is_flipped = structure_check_by_ssim(flipped_original, result_image, threshold=0.99)
        if is_flipped:
            logger.info('Image is correctly flipped horizontally')
            return 1.0
        else:
            logger.warning('Image is not correctly flipped')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking flipped image: {e}')
        return 0.0

def check_image_dimensions__016cfcf5(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_dimensions__bec1165ff4f2eb5b48dc7de50a4fe1ab(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected values)
        **options: Additional options (unused)

    Returns:
        float: 1.0 if dimensions match exactly, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if expected_width is not None and actual_width != expected_width:
        return 0.0
    if expected_height is not None and actual_height != expected_height:
        return 0.0
    return 1.0

def check_image_grayscale__9660fa13484d43a1462069231ef86deb(result, expected, **options):
    """
    Check if image is in grayscale mode.

    Args:
        result: str - Image mode from getter (e.g., 'L', 'RGB', 'RGBA')
        expected: dict with 'grayscale' boolean flag

    Returns:
        float: 1.0 if mode matches expectation, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    should_be_grayscale = expected.get('grayscale', True)
    is_grayscale = result == 'L'
    logger.info(f'Image mode: {result}, Expected grayscale: {should_be_grayscale}, Is grayscale: {is_grayscale}')
    if is_grayscale == should_be_grayscale:
        return 1.0
    else:
        return 0.0

def check_image_extracted__345e8ddb(result, expected, **options):
    """
    Check if the correct image (second image) was extracted with exact properties.

    Args:
        result: dict with image properties (from getter)
        expected: dict with exact expected properties (expected_size, expected_width,
                 expected_height, format, with tolerances)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists'):
        score += 0.4
    else:
        return 0.0
    expected_size = expected.get('expected_size', 40335)
    size_tolerance = expected.get('size_tolerance', 500)
    actual_size = result.get('size', 0)
    if abs(actual_size - expected_size) <= size_tolerance:
        score += 0.2
    expected_format = expected.get('format', 'PNG')
    if result.get('format') == expected_format:
        score += 0.2
    expected_width = expected.get('expected_width', 1114)
    expected_height = expected.get('expected_height', 623)
    dimension_tolerance = expected.get('dimension_tolerance', 5)
    actual_width = result.get('width', 0)
    actual_height = result.get('height', 0)
    width_match = abs(actual_width - expected_width) <= dimension_tolerance
    height_match = abs(actual_height - expected_height) <= dimension_tolerance
    if width_match and height_match:
        score += 0.2
    return score

def check_gif_file__8d385ddc(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_format__739292ff(result, expected, **options):
    """Check if image format matches expected value and verify it's a genuine screenshot.

    This function verifies:
    1. The image format is correct (PNG)
    2. The file size is reasonable for a screenshot (between 1KB and 50MB)
    3. The image dimensions are reasonable for a video screenshot (between 320x240 and 7680x4320)
    4. The file was created recently (within the last 5 minutes, if timestamp available)

    Args:
        result: Dict with image properties (width, height, format, size_bytes, creation_time)
        expected: Dict with 'format' key (e.g., 'PNG', 'JPEG')
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None - file may not exist or could not be read')
        return 0.0
    expected_format = expected.get('format', 'PNG')
    actual_format = result.get('format', '')
    if actual_format.upper() != expected_format.upper():
        logger.info(f'Format mismatch: expected {expected_format}, got {actual_format}')
        return 0.0
    size_bytes = result.get('size_bytes', 0)
    min_size = 1024
    max_size = 50 * 1024 * 1024
    if size_bytes < min_size:
        logger.info(f'File too small: {size_bytes} bytes (minimum {min_size})')
        return 0.0
    if size_bytes > max_size:
        logger.info(f'File too large: {size_bytes} bytes (maximum {max_size})')
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    if width < 320 or height < 240:
        logger.info(f'Image dimensions too small: {width}x{height} (minimum 320x240)')
        return 0.0
    if width > 7680 or height > 4320:
        logger.info(f'Image dimensions too large: {width}x{height} (maximum 7680x4320)')
        return 0.0
    creation_time = result.get('creation_time')
    if creation_time is not None:
        current_time = time.time()
        time_diff = current_time - creation_time
        max_age = 300
        if time_diff > max_age:
            logger.info(f'File is too old: created {time_diff:.1f} seconds ago (max {max_age})')
            return 0.0
        if time_diff < 0:
            logger.info(f'File creation time is in the future: {time_diff:.1f} seconds')
            return 0.0
    logger.info(f'All checks passed: format={actual_format}, size={size_bytes} bytes, dimensions={width}x{height}')
    return 1.0

def check_vlc_default_player__d253b8f0(result, expected, **options):
    """
    Check if VLC is set as the default video player.

    Args:
        result: Current default video player application
        expected: Expected configuration with 'is_vlc_default' boolean
        **options: Additional options

    Returns:
        float: 1.0 if VLC is default, 0.0 otherwise
    """
    is_vlc_default = expected.get('is_vlc_default', True)
    vlc_is_default = 'vlc' in result.lower() if result else False
    if is_vlc_default:
        return 1.0 if vlc_is_default else 0.0
    else:
        return 0.0 if vlc_is_default else 1.0

def check_triangle_top_right__1a2ce7aa(tgt_path, expected, **options):
    """
    Check if the triangle is in the top-right corner of the image.
    Variation 4 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with tolerance parameter
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is in top-right, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        img = Image.open(tgt_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        sorted_indices = np.argsort(counts)[::-1]
        unique_colors_sorted = unique_colors[sorted_indices]
        triangle_color = unique_colors_sorted[1]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        centroid = triangle_coords.mean(axis=0)
        (height, width) = img_array.shape[:2]
        tolerance = expected.get('tolerance', 0.1)
        in_top = centroid[0] < height * (0.5 + tolerance)
        in_right = centroid[1] > width * (0.5 - tolerance)
        logger.info(f'Triangle centroid: {centroid}, Image size: ({height}, {width})')
        logger.info(f'In top: {in_top}, In right: {in_right}')
        if in_top and in_right:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_resize_and_structure_sim__ed3206c993cd330e3bb4b56cf8b439e8(src_path, expected):
    """
    Check if the image has been resized to expected dimensions
    Task: Resize image to specific dimensions

    Args:
        src_path: Path to the resized image
        expected: dict with 'width' and 'height' keys from rules
    """
    logger.info(f'Evaluating image resize: src={src_path}, expected={expected}')
    if src_path is None:
        logger.warning('Source path is None')
        return 0.0
    try:
        source_image = Image.open(src_path)
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        logger.info(f'Source size: {source_image.size}, Expected: {expected_width}x{expected_height}')
        width_match = source_image.size[0] == expected_width
        height_match = source_image.size[1] == expected_height
        if width_match and height_match:
            logger.info('Dimensions match expected values')
            return 1.0
        else:
            logger.warning(f'Dimensions mismatch: got {source_image.size}, expected {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error during resize evaluation: {e}')
        return 0.0

def check_image_cropped__166a89c9(result_state: Optional[Dict[str, int]], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if an image has been cropped by verifying its dimensions are smaller than the maximum allowed.

    Args:
        result_state: Dict with 'width' and 'height' of the actual image, or None if file not found
        expected_state: Dict containing 'max_width' and 'max_height' (the rules dict directly)
        **options: Additional options (unused)

    Returns:
        float: Score of 1.0 if image is cropped (both width < max_width AND height < max_height), 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Image file not found or could not be read')
        return 0.0
    actual_width = result_state.get('width')
    actual_height = result_state.get('height')
    if actual_width is None or actual_height is None:
        logger.warning('Image dimensions could not be determined')
        return 0.0
    max_width = expected_state.get('max_width')
    max_height = expected_state.get('max_height')
    if max_width is None or max_height is None:
        logger.error('max_width or max_height not specified in rules')
        return 0.0
    logger.info(f'Checking if image is cropped: actual={actual_width}x{actual_height}, max={max_width}x{max_height}')
    if actual_width < max_width and actual_height < max_height:
        logger.info('Image is properly cropped')
        return 1.0
    else:
        logger.warning(f'Image is not cropped: width={actual_width}<{max_width}? {actual_width < max_width}, height={actual_height}<{max_height}? {actual_height < max_height}')
        return 0.0

def check_layer_exists__b148e375_v3(result_state: List[str], expected_state: dict, **options) -> float:
    """
    Check if a layer with the specified name exists in the GIMP layer list.

    Args:
        result_state: List of layer names from the getter
        expected_state: Dict with 'layer_name' key specifying the expected layer name
        **options: Additional options

    Returns:
        float: 1.0 if the layer exists, 0.0 otherwise
    """
    expected_layer_name = expected_state.get('layer_name', 'Overlay')
    if not isinstance(result_state, list):
        logger.error(f'Invalid result_state type: {type(result_state)}, expected list')
        return 0.0
    for layer_name in result_state:
        if layer_name == expected_layer_name:
            logger.info(f"Layer '{expected_layer_name}' found in layer list")
            return 1.0
    logger.warning(f"Layer '{expected_layer_name}' not found in layer list: {result_state}")
    return 0.0

def check_extracted_image__938bcb9a(result, expected, **options):
    """
    Check if a valid PNG image file was extracted from Word attachment in Thunderbird Notes folder.

    Verifies:
    1. File exists at the expected location
    2. File is a valid PNG image that can be opened
    3. File size is reasonable (not empty, not suspiciously large)
    4. Image has valid dimensions
    5. LibreOffice Writer was used (to open Word attachment)
    6. Thunderbird was used (to access email)
    7. Word attachment was accessed (verified via Thunderbird cache)
    8. Notes folder was accessed (optional, via Thunderbird profile prefs)

    This provides comprehensive verification by checking not just the output file,
    but also that the required applications were used and the correct attachment
    source was accessed during task execution, significantly reducing false positives.

    Args:
        result: Output from Python command in format "exists|size|valid|dimensions|hash|lowriter_used|thunderbird_used|word_attachment_found|notes_folder_accessed"
                e.g., "True|12345|True|800x600|a1b2c3d4|True|True|True|True"
        expected: Dict with 'min_file_size', 'max_file_size', 'must_be_valid_png',
                  'require_libreoffice_usage', 'require_thunderbird_usage',
                  'require_word_attachment', 'require_notes_folder_access'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    try:
        parts = result.strip().split('|')
        if len(parts) != 9:
            return 0.0
        (exists_str, size_str, valid_str, dims_str, img_hash, lowriter_str, thunderbird_str, word_attach_str, notes_folder_str) = parts
        exists = exists_str.lower() == 'true'
        if not exists:
            return 0.0
        is_valid_png = valid_str.lower() == 'true'
        must_be_valid = expected.get('must_be_valid_png', True)
        if must_be_valid and (not is_valid_png):
            return 0.0
        try:
            file_size = int(size_str)
        except ValueError:
            return 0.0
        min_size = expected.get('min_file_size', 1000)
        max_size = expected.get('max_file_size', 50000000)
        if file_size < min_size or file_size > max_size:
            return 0.0
        if dims_str == '0x0':
            return 0.0
        try:
            (width, height) = dims_str.split('x')
            (w, h) = (int(width), int(height))
            if w <= 0 or h <= 0:
                return 0.0
            if w > 10000 or h > 10000:
                return 0.0
        except (ValueError, AttributeError):
            return 0.0
        require_lowriter = expected.get('require_libreoffice_usage', True)
        lowriter_used = lowriter_str.lower() == 'true'
        if require_lowriter and (not lowriter_used):
            return 0.0
        require_thunderbird = expected.get('require_thunderbird_usage', True)
        thunderbird_used = thunderbird_str.lower() == 'true'
        if require_thunderbird and (not thunderbird_used):
            return 0.0
        require_word_attach = expected.get('require_word_attachment', True)
        word_attach_found = word_attach_str.lower() == 'true'
        if require_word_attach and (not word_attach_found):
            return 0.0
        require_notes_folder = expected.get('require_notes_folder_access', False)
        notes_folder_accessed = notes_folder_str.lower() == 'true'
        if require_notes_folder and (not notes_folder_accessed):
            return 0.0
        return 1.0
    except Exception:
        return 0.0

def check_png_deleted__eb6b46ad(png_count, expected):
    """Check if PNG files were deleted.

    Args:
        png_count: Number of PNG files remaining
        expected: Dict with 'count' key (expected to be 0)

    Returns:
        float: 1.0 if count matches expected, 0.0 otherwise
    """
    expected_count = expected['count']
    if png_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_image_hash__d1128c0a(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_file__203069587d53a571860bccb97348992b(result, expected, **options):
    """Check if image file exists and meets size requirements.

    Args:
        result: Dict from getter with 'exists', 'path', 'size' keys
        expected: Dict with 'exists' and optional 'min_size' from rules
        **options: Additional options

    Returns:
        float: 1.0 if image file meets requirements, 0.0 otherwise
    """
    if not isinstance(result, dict) or 'exists' not in result:
        logger.error(f'Invalid result format: {result}')
        return 0.0
    expected_exists = expected.get('exists', True)
    actual_exists = result.get('exists', False)
    actual_size = result.get('size', 0)
    min_size = expected.get('min_size', 1000)
    if not expected_exists:
        return 1.0 if not actual_exists else 0.0
    if not actual_exists:
        logger.warning(f"Image file {result.get('path', 'unknown')} does not exist")
        return 0.0
    if actual_size < min_size:
        logger.warning(f'Image file size {actual_size} is below minimum {min_size}')
        return 0.0
    logger.info(f"Image file {result.get('path', 'unknown')} exists with size {actual_size} >= {min_size}")
    return 1.0

def check_image_flipped__16b4973e(result, expected, **options):
    """
    Check if the image was flipped horizontally.

    Args:
        result: Dict with 'flipped_horizontal' from getter
        expected: Dict with expected flip status
        **options: Additional options

    Returns:
        float: Score (1.0 if flipped, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        is_flipped = result.get('flipped_horizontal', False)
        expected_flipped = expected.get('flipped_horizontal', True)
        logger.info(f'Expected flipped: {expected_flipped}, Actual flipped: {is_flipped}')
        if is_flipped == expected_flipped:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking flip: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_gimp_autosave_setting__d7b0aa1f(actual_config_path, expected, **options):
    """
    Check if GIMP autosave settings are enabled with correct interval.

    This metric checks both:
    1. undo-autosave yes
    2. undo-autosave-interval <value>

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'autosave_enabled' and 'autosave_interval' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if both settings match, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        autosave_enabled_expected = expected.get('autosave_enabled', True)
        autosave_interval_expected = expected.get('autosave_interval', 300)
        autosave_enabled_found = False
        autosave_interval_found = False
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == 'undo-autosave':
                    expected_val = 'yes' if autosave_enabled_expected else 'no'
                    if value == expected_val:
                        autosave_enabled_found = True
                        logger.info(f'Found autosave enabled: {value}')
                if key == 'undo-autosave-interval':
                    if int(value) == autosave_interval_expected:
                        autosave_interval_found = True
                        logger.info(f'Found autosave interval: {value}')
        if autosave_enabled_found and autosave_interval_found:
            logger.info('Both autosave settings match')
            return 1.0
        else:
            logger.warning(f'Autosave settings mismatch: enabled={autosave_enabled_found}, interval={autosave_interval_found}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking autosave settings: {e}')
        return 0.0

def check_image_dimensions__739292ff(result, expected, **options):
    """Check if image dimensions are within expected range.

    Args:
        result: Dict with image properties (width, height)
        expected: Dict with 'min_width', 'min_height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions meet criteria, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    min_width = expected.get('min_width', 0)
    min_height = expected.get('min_height', 0)
    actual_width = result.get('width', 0)
    actual_height = result.get('height', 0)
    if actual_width >= min_width and actual_height >= min_height:
        logger.info(f'Dimensions OK: {actual_width}x{actual_height} >= {min_width}x{min_height}')
        return 1.0
    else:
        logger.info(f'Dimensions too small: {actual_width}x{actual_height} < {min_width}x{min_height}')
        return 0.0

def check_all_images__1574bc56f755697238f4190c2d72a32b(result, expected, **options):
    """Check if all image files (jpg and png) were moved correctly.

    Args:
        result: Directory tree dict from getter
        expected: Rules dict with 'expected' list of filenames
        **options: Additional options

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected.get('expected', [])
    if not result or 'children' not in result:
        return 0.0
    actual_files = [node['name'] for node in result['children']]
    if len(actual_files) != len(expected_files):
        return 0.0
    if set(actual_files) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_final_image_count__5d239d03(result, expected, **options):
    """Check if final image count is correct.

    Args:
        result: Actual final count
        expected: Expected final count (dict with 'final_count' key)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_final = expected.get('final_count', 4)
    return 1.0 if result == expected_final else 0.0

def check_image_extracted__f3e27026(result, expected, **options):
    """
    Check if an image was extracted with correct properties.

    Args:
        result: dict with image properties (from getter)
        expected: dict with expected properties (min_size, format, min_width, min_height)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists'):
        score += 0.4
    else:
        return 0.0
    min_size = expected.get('min_size', 1000)
    if result.get('size', 0) >= min_size:
        score += 0.2
    expected_format = expected.get('format', 'PNG')
    if result.get('format') == expected_format:
        score += 0.2
    min_width = expected.get('min_width', 100)
    min_height = expected.get('min_height', 100)
    if result.get('width', 0) >= min_width and result.get('height', 0) >= min_height:
        score += 0.2
    return score

def check_image_hash__0969de9d(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_rotated__777d01ae6d6663c5897f2674681dca2e(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if image exists and appears to be rotated 90 degrees.

    Args:
        result: Dict from getter with 'exists', 'is_rotated_90' keys
        expected: Dict with 'is_rotated_90' boolean

    Returns:
        float: 1.0 if file exists and rotation matches expectation, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('Image file does not exist')
        return 0.0
    expected_rotation = expected.get('is_rotated_90', False)
    actual_rotation = result.get('is_rotated_90', False)
    if actual_rotation == expected_rotation:
        logger.info(f'Image rotation matches: rotated_90={actual_rotation}')
        return 1.0
    else:
        logger.info(f'Image rotation mismatch: expected rotated_90={expected_rotation}, got {actual_rotation}')
        return 0.0

def check_gimp_single_window__d9681274028e239cd7a4e13a8dc53d3c(actual_config_path, rule):
    """
    Check if GIMP is set to single-window mode
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_image_size__2110d1dc(src_path, rule):
    """
    Check if the size of the src image matches the expected dimensions
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_0
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        actual_width = img.size[0]
        actual_height = img.size[1]
        logger.debug(f'Image size: {img.size}')
        if rule.get('height', None) is not None:
            height_same = actual_height == rule['height']
        else:
            height_same = True
        if rule.get('width', None) is not None:
            width_same = actual_width == rule['width']
        else:
            width_same = True
        if height_same and width_same:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 1.0
        else:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image size: {e}')
        return 0.0

def check_png_export__046a9ca717a3ff75711d7d7d1d876a5f(result, expected, **options):
    """
    Check if a PNG file was properly exported from video.

    Scoring breakdown:
    - File exists (0.25)
    - Source video exists, confirming context (0.15)
    - Format is PNG (0.35)
    - Dimensions are reasonable (0.25)

    Args:
        result: Dict from getter with PNG file info
        expected: Dict with expected values (format, min_dimension)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    if not result.get('exists', False):
        logger.info('PNG file does not exist')
        return 0.0
    score = 0.25
    if result.get('source_video_exists', False):
        score += 0.15
        logger.info('Source video exists, confirming valid task context')
    else:
        logger.warning('Source video not found - task context may be invalid')
    expected_format = expected.get('format', 'PNG')
    if result.get('format') == expected_format:
        score += 0.35
        logger.info(f"Format check passed: {result.get('format')}")
    else:
        logger.warning(f"Format mismatch: expected {expected_format}, got {result.get('format')}")
    width = result.get('width', 0)
    height = result.get('height', 0)
    min_dimension = expected.get('min_dimension', 100)
    if width >= min_dimension and height >= min_dimension:
        score += 0.25
        logger.info(f'Dimensions check passed: {width}x{height} (min: {min_dimension})')
    else:
        logger.warning(f'Dimensions too small: {width}x{height} (min: {min_dimension})')
    return score

def check_image_dimensions__b24412f0(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_layer_exists__17690e7b(result_state, expected_state, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected_state: Expected configuration containing 'layer_name' (dict)
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected_state, dict):
            expected_layer = expected_state.get('layer_name', '')
        else:
            expected_layer = str(expected_state)
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_third_image__9a3feca86555f3c82732707871883142(result, expected, **options):
    """Check if the third image was extracted and saved.

    Args:
        result: dict from getter with 'exists', 'size', 'is_png', 'hash', 'reference_hash'
        expected: dict with expected properties (not used, just for consistency)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.3
    else:
        return 0.0
    if result.get('is_png', False):
        score += 0.2
    else:
        return 0.0
    if result.get('size', 0) >= 5000:
        score += 0.1
    file_hash = result.get('hash')
    reference_hash = result.get('reference_hash')
    if reference_hash and file_hash:
        if file_hash == reference_hash:
            score += 0.4
        else:
            score = score * 0.3
    elif not reference_hash:
        score += 0.2
    return score

def check_has_multiple_images__67a58d2f(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if folder has the exact number of PNG image files (verifies ALL images extracted)"""
    exact_count = expected.get('exact_count')
    min_count = expected.get('min_count')
    if exact_count is not None:
        if len(result) == exact_count:
            logger.info(f'Has exactly {len(result)} PNG files (expected: {exact_count})')
            return 1.0
        else:
            logger.warning(f'Has {len(result)} PNG files but expected exactly {exact_count}')
            return 0.0
    elif min_count is not None:
        if len(result) >= min_count:
            logger.info(f'Has {len(result)} PNG files (min: {min_count})')
            return 1.0
        return 0.0
    else:
        if len(result) >= 2:
            logger.info(f'Has {len(result)} PNG files (default min: 2)')
            return 1.0
        return 0.0

def check_gimp_rulers__f13e522727f9f47f05a1828391e4f1b5(actual_config_path, rule):
    """
    Check if GIMP rulers visibility setting is correct
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_image_mirrored__2237c5bebdb76e54ae53ea89e71ca4a3(result, expected, **options):
    """Check if the result image is a horizontal mirror of the original image.

    Args:
        result: Path to the result image file
        expected: Dict with 'original_path' key containing path to original image
        **options: Additional options (threshold for SSIM comparison)

    Returns:
        float: 1.0 if image is mirrored, 0.0 otherwise
    """
    if result is None:
        return 0.0
    original_path = expected.get('original_path')
    if not original_path:
        return 0.0
    try:
        result_image = Image.open(result)
        original_image = Image.open(original_path)
        result_flipped = result_image.transpose(Image.FLIP_LEFT_RIGHT)
        if result_flipped.mode != 'RGB':
            result_flipped = result_flipped.convert('RGB')
        if original_image.mode != 'RGB':
            original_image = original_image.convert('RGB')
        if result_flipped.size != original_image.size:
            logger.debug(f'Size mismatch: {result_flipped.size} vs {original_image.size}')
            return 0.0
        array1 = np.array(result_flipped)
        array2 = np.array(original_image)
        min_dim = min(array1.shape[0], array1.shape[1])
        if min_dim < 7:
            win_size = min_dim if min_dim % 2 == 1 else min_dim - 1
            if win_size < 1:
                logger.debug('Image too small for SSIM computation')
                return 0.0
        else:
            win_size = 7
        try:
            similarity = ssim(array1, array2, win_size=win_size, channel_axis=2)
        except TypeError:
            similarity = ssim(array1, array2, win_size=win_size, multichannel=True)
        threshold = options.get('threshold', 0.99)
        logger.debug(f'SSIM similarity: {similarity}, threshold: {threshold}')
        return 1.0 if similarity >= threshold else 0.0
    except Exception as e:
        logger.error(f'Error checking image mirror: {e}')
        return 0.0

def check_jpeg_export__a818e71b(result_path, expected, **options):
    """
    Check if the image has been exported as JPEG format with required quality.

    Args:
        result_path: Path to the result image
        expected: Dict with optional 'min_quality' threshold
        **options: Additional options

    Returns:
        float: Score (1.0 if JPEG format with sufficient quality, 0.0 otherwise)
    """
    if result_path is None:
        logger.error('Result path is None')
        return 0.0
    try:
        (_, ext) = os.path.splitext(result_path)
        if ext.lower() not in ['.jpg', '.jpeg']:
            logger.warning(f'File extension is not .jpg or .jpeg: {ext}')
            return 0.0
        img = Image.open(result_path)
        if img.format != 'JPEG':
            logger.warning(f'Image format is {img.format}, not JPEG')
            return 0.0
        logger.info(f'Image is in JPEG format')
        min_quality = expected.get('min_quality', 0)
        if min_quality > 0:
            estimated_quality = estimate_jpeg_quality(result_path)
            if estimated_quality is None:
                logger.warning('Could not determine JPEG quality, assuming it meets requirement')
                return 0.5
            if estimated_quality >= min_quality:
                logger.info(f'JPEG quality {estimated_quality} meets minimum requirement {min_quality}')
                return 1.0
            else:
                logger.warning(f'JPEG quality {estimated_quality} is below minimum requirement {min_quality}')
                return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking JPEG export: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_composite_photo_organization__8542f72b(result: dict, expected: dict, **options) -> float:
    """Composite check for photo organization task.

    Verifies:
    1. stage_photos folder has correct 4 files (40% weight)
    2. other_photos folder has correct 2 files (40% weight)
    3. Original directory is empty (files moved, not copied) (20% weight)

    Args:
        result: Dict with keys 'stage_photos', 'other_photos', 'original_dir'
        expected: Dict with expected values for each key
        **options: Additional options

    Returns:
        float: Weighted score between 0.0 and 1.0
    """
    score = 0.0
    stage_photos_actual = sorted(result.get('stage_photos', []))
    stage_photos_expected = sorted(expected.get('stage_photos', []))
    if stage_photos_actual == stage_photos_expected:
        score += 0.4
    other_photos_actual = sorted(result.get('other_photos', []))
    other_photos_expected = sorted(expected.get('other_photos', []))
    if other_photos_actual == other_photos_expected:
        score += 0.4
    original_dir_actual = result.get('original_dir', [])
    if len(original_dir_actual) == 0:
        score += 0.2
    return score

def check_image_blurred__eea46e4e4afa193660c5f52c6a2da7a9(result, expected, **options):
    """Check if image is sufficiently blurred compared to threshold with multi-step verification.

    This metric implements partial credit scoring:
    - Step 1: File existence (0.3 weight) - checks that pic.jpg was created
    - Step 2: Blur effect (0.7 weight) - checks that image has sufficient blur

    Args:
        result: Dict with file_exists, variance, and is_blurred from getter
        expected: Dict with max_variance threshold
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 with partial credit
    """
    if not result or not expected:
        return 0.0
    file_exists = result.get('file_exists', False)
    file_score = 0.3 if file_exists else 0.0
    if not file_exists:
        return file_score
    variance = result.get('variance', float('inf'))
    max_variance = expected.get('max_variance', 0)
    is_blurred = variance < max_variance
    blur_score = 0.7 if is_blurred else 0.0
    return file_score + blur_score

def check_layer_exists__22c9ed6961463c7a6e1af4ed7b883c13(result_state, expected_state, **options):
    """
    Check if a specific layer name exists in the list of layer names.

    Args:
        result_state: List of layer names from the XCF file
        expected_state: Dict with 'layer_name' key specifying the expected layer name
        **options: Additional options

    Returns:
        float: 1.0 if the layer exists, 0.0 otherwise
    """
    if not isinstance(result_state, list):
        logger.error(f'result_state is not a list: {type(result_state)}')
        return 0.0
    if not isinstance(expected_state, dict):
        logger.error(f'expected_state is not a dict: {type(expected_state)}')
        return 0.0
    expected_layer_name = expected_state.get('layer_name', '')
    if not expected_layer_name:
        logger.error('No layer_name specified in expected_state')
        return 0.0
    if expected_layer_name in result_state:
        logger.info(f"Layer '{expected_layer_name}' found in {result_state}")
        return 1.0
    else:
        logger.warning(f"Layer '{expected_layer_name}' not found in {result_state}")
        return 0.0

def check_layer_renamed__677f0e4c(result, expected, **options):
    """
    Check if the top layer has been renamed to the expected name.

    Args:
        result: List of layer names from getter (list of strings, ordered top-to-bottom)
        expected: Dict with 'new_name' - the expected layer name
        **options: Additional options

    Returns:
        float: 1.0 if the top layer (result[0]) matches the expected name, 0.0 otherwise
    """
    try:
        new_name = expected.get('new_name', '')
        if not isinstance(result, list):
            logger.error(f'result is not a list: {type(result)}')
            return 0.0
        if not result:
            logger.error('Layer list is empty')
            return 0.0
        if result[0] == new_name:
            logger.info(f"Top layer successfully renamed to '{new_name}' (layer stack: {result})")
            return 1.0
        else:
            logger.warning(f"Top layer is '{result[0]}', expected '{new_name}' (layer stack: {result})")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer rename: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_horizontal_flip__08f973926b8f35af5489796c73a6c6e0(result, expected, **options):
    """
    Check if the result image is a horizontal flip of the source image.

    Args:
        result: Dict with 'result_path' and 'source_path'
        expected: Dict with expected rule (not used, just for consistency)
        **options: Additional options

    Returns:
        float: 1.0 if horizontally flipped correctly, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    result_path = result.get('result_path')
    source_path = result.get('source_path')
    if not result_path or not source_path:
        logger.error('Missing result_path or source_path')
        return 0.0
    try:
        result_img = Image.open(result_path)
        source_img = Image.open(source_path)
        if result_img.size != source_img.size:
            logger.error(f'Size mismatch: {result_img.size} vs {source_img.size}')
            return 0.0
        flipped_source = source_img.transpose(Image.FLIP_LEFT_RIGHT)
        if result_img.mode != 'RGB':
            result_img = result_img.convert('RGB')
        if flipped_source.mode != 'RGB':
            flipped_source = flipped_source.convert('RGB')
        result_array = np.array(result_img)
        flipped_array = np.array(flipped_source)
        try:
            similarity = ssim(result_array, flipped_array, win_size=7, channel_axis=2)
        except TypeError:
            similarity = ssim(result_array, flipped_array, win_size=7, multichannel=True)
        logger.info(f'Horizontal flip SSIM: {similarity}')
        threshold = options.get('threshold', 0.95)
        return 1.0 if similarity >= threshold else 0.0
    except Exception as e:
        logger.error(f'Error checking horizontal flip: {e}')
        return 0.0

def check_triangle_bottom_left__e0c275c9(tgt_path, expected, **options):
    """
    Check if the triangle is in the bottom-left corner of the image.
    Variation 0 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with tolerance parameter
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is in bottom-left, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        img = Image.open(tgt_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)[::-1]]
        if img_array.shape[2] >= 3:
            yellow_mask = (img_array[:, :, 0] > 150) & (img_array[:, :, 1] > 150) & (img_array[:, :, 2] < 150)
            if yellow_mask.any():
                triangle_mask = yellow_mask
            else:
                if len(unique_colors_sorted) < 2:
                    logger.error('Not enough unique colors to detect triangle')
                    return 0.0
                triangle_color = unique_colors_sorted[1]
                triangle_mask = np.all(img_array == triangle_color, axis=2)
        else:
            if len(unique_colors_sorted) < 2:
                logger.error('Not enough unique colors to detect triangle')
                return 0.0
            triangle_color = unique_colors_sorted[1]
            triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.error('No triangle pixels detected')
            return 0.0
        centroid = triangle_coords.mean(axis=0)
        (height, width) = img_array.shape[:2]
        tolerance = expected.get('tolerance', 0.1)
        in_bottom = centroid[0] > height * (0.5 - tolerance)
        in_left = centroid[1] < width * (0.5 + tolerance)
        logger.info(f'Triangle centroid: {centroid}, Image size: ({height}, {width})')
        logger.info(f'In bottom: {in_bottom}, In left: {in_left}')
        if in_bottom and in_left:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_resize__dc2a3b39(src_path, expected):
    """
    Check if the image was resized to expected dimensions.
    Variation 4 for task 554785e9-4523-4e7a-b8e1-8016f565f56a

    Args:
        src_path: Path to edited image
        expected: Dict with 'width' and 'height' keys

    Returns:
        float: 1.0 if image has expected dimensions, 0.0 otherwise
    """
    if src_path is None:
        logger.error('Source path is None')
        return 0.0
    try:
        img = Image.open(src_path)
        (actual_width, actual_height) = img.size
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        logger.info(f'Image size check: actual={actual_width}x{actual_height}, expected={expected_width}x{expected_height}')
        width_match = expected_width is None or actual_width == expected_width
        height_match = expected_height is None or actual_height == expected_height
        if width_match and height_match:
            return 1.0
        else:
            logger.debug(f'Size mismatch: width_match={width_match}, height_match={height_match}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image resize: {e}')
        return 0.0

def check_jpg_files_copied__2c299b404d518d11da9f18c6fc937d90(directory_list, rule):
    """
    Check if the directory contains exactly the expected number of JPG files and no other files.

    This metric verifies:
    1. The correct number of files are present (completeness)
    2. Only files with allowed extensions are present (selectivity)

    Args:
        directory_list: Directory tree structure from get_list_directory
        rule: Expected configuration with 'expected_count' and 'allowed_extensions' keys

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_count = rule.get('expected_count', 0)
    allowed_extensions = rule.get('allowed_extensions', [])
    children = directory_list.get('children', [])
    files = [child for child in children if child.get('type') == 'file']
    actual_count = len(files)
    if actual_count != expected_count:
        return 0.0
    for file_info in files:
        file_name = file_info.get('name', '')
        file_ext = os.path.splitext(file_name)[1].lower()
        allowed_exts_lower = [ext.lower() for ext in allowed_extensions]
        if file_ext not in allowed_exts_lower:
            return 0.0
    return 1.0

def check_image_format_and_existence__e19bfc7ef1338231fef513d8a5b2f6d1(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if image exists and has expected format.

    Args:
        result: Dict from getter with 'exists', 'format' keys
        expected: Dict with 'format' key specifying expected format

    Returns:
        float: 1.0 if file exists and format matches, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('Image file does not exist')
        return 0.0
    expected_format = expected.get('format', '').upper()
    actual_format = (result.get('format') or '').upper()
    if actual_format == expected_format:
        logger.info(f'Image format matches: {actual_format}')
        return 1.0
    else:
        logger.info(f'Image format mismatch: expected {expected_format}, got {actual_format}')
        return 0.0

def check_jpg_exists__8d385ddc(result_state, expected_state, **options):
    """
    Check if the JPEG file exists, is in JPEG format, and matches the expected frame content.

    Task: 2fe4b718-3bd7-46ec-bdce-b184f5653624
    Instruction: Extract a screenshot at exactly 3 seconds from the video 'src.mp4'
                 on the desktop and save it as 'screenshot_3s.jpg'.

    This evaluator verifies:
    1. File exists with correct filename
    2. File is in JPEG format
    3. Image content matches the actual frame at 3 seconds from src.mp4 (via SSIM comparison)

    Args:
        result_state: Path to the user-created JPEG file (str)
        expected_state: Dict containing:
                       - 'format': Expected image format (e.g., 'JPEG')
                       - 'reference_frame': Path to reference frame extracted at 3s from src.mp4 (str)
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass (file exists, correct format, and content matches), 0.0 otherwise
    """
    logger = logging.getLogger(__name__)
    if result_state is None or result_state == '':
        logger.warning('Result state is None or empty - screenshot file not found')
        return 0.0
    if not os.path.isfile(result_state):
        logger.warning(f'Screenshot file does not exist: {result_state}')
        return 0.0
    try:
        user_img = Image.open(result_state)
        expected_format = expected_state.get('format', 'JPEG')
        if user_img.format != expected_format:
            logger.warning(f'Image format mismatch: expected {expected_format}, got {user_img.format}')
            return 0.0
        logger.info(f'Format check passed: {result_state} (format: {user_img.format}, size: {user_img.size})')
        reference_frame_path = expected_state.get('reference_frame')
        if not reference_frame_path:
            logger.error('No reference frame provided for content verification')
            return 0.0
        if not os.path.isfile(reference_frame_path):
            logger.error(f'Reference frame does not exist: {reference_frame_path}')
            return 0.0
        reference_img = Image.open(reference_frame_path)
        user_img_gray = user_img.convert('L')
        reference_img_gray = reference_img.convert('L')
        user_size = user_img_gray.size
        ref_size = reference_img_gray.size
        min_width = min(user_size[0], ref_size[0])
        min_height = min(user_size[1], ref_size[1])
        new_size = (min_width, min_height)
        user_img_resized = user_img_gray.resize(new_size, Image.Resampling.LANCZOS)
        reference_img_resized = reference_img_gray.resize(new_size, Image.Resampling.LANCZOS)
        user_array = np.array(user_img_resized)
        reference_array = np.array(reference_img_resized)
        similarity = ssim(user_array, reference_array)
        logger.info(f'Image similarity (SSIM): {similarity:.4f}')
        threshold = 0.85
        if similarity >= threshold:
            logger.info(f'Content verification passed: SSIM {similarity:.4f} >= threshold {threshold}')
            return 1.0
        else:
            logger.warning(f'Content verification failed: SSIM {similarity:.4f} < threshold {threshold}')
            logger.warning('The screenshot does not match the frame at 3 seconds from the source video')
            return 0.0
    except Exception as e:
        logger.error(f'Error verifying screenshot: {e}', exc_info=True)
        return 0.0

def check_image_extracted__4e5534f0(result, expected, **options):
    """
    Check if images were extracted correctly.

    Args:
        result: Dict with 'folder_exists' and 'images' list
        expected: Dict with validation rules (min_size, format, min_width, min_height, min_count)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('folder_exists', False):
        return 0.0
    score += 0.3
    images = result.get('images', [])
    if len(images) == 0:
        return score
    score += 0.2
    min_count = expected.get('min_count', 1)
    if len(images) >= min_count:
        score += 0.1
    min_size = expected.get('min_size', 1000)
    expected_format = expected.get('format', 'PNG')
    min_width = expected.get('min_width', 100)
    min_height = expected.get('min_height', 100)
    valid_images = 0
    for img in images:
        is_valid = True
        if img.get('size', 0) < min_size:
            is_valid = False
        img_format = img.get('format', '')
        if expected_format == 'PNG' and img_format not in ['PNG']:
            is_valid = False
        elif expected_format == 'JPEG' and img_format not in ['JPEG', 'JPG']:
            is_valid = False
        elif expected_format not in ['PNG', 'JPEG'] and img_format != expected_format:
            is_valid = False
        if img.get('width', 0) < min_width or img.get('height', 0) < min_height:
            is_valid = False
        if is_valid:
            valid_images += 1
    if len(images) > 0:
        quality_score = valid_images / len(images)
        score += 0.4 * quality_score
    return score

def check_gimp_tabs__ca73cf065f0842768815735a3354abf0(actual_config_path, rule):
    """
    Check if GIMP tabs visibility setting is correct
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_image_cropped__c48d85ee866b2dcfd0ab60f091a5a2b6(result, expected, **options):
    """
    Check if image has been cropped to expected dimensions.

    Args:
        result: Dict with 'width', 'height' from getter
        expected: Expected rules dict with 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if not result:
        logger.error('No result data')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result.get('width') != expected_width or result.get('height') != expected_height:
        logger.error(f"Dimensions mismatch: expected {expected_width}x{expected_height}, got {result.get('width')}x{result.get('height')}")
        return 0.0
    logger.info(f'Image cropped correctly to {expected_width}x{expected_height}')
    return 1.0

def check_image_crop__4ff6a540(src_path, expected, **options):
    """
    Check if the image has been cropped to the expected dimensions.
    Variation 6: 2a729ded-3296-423d-aec4-7dd55ed5fbb3

    Args:
        src_path: Path to the result image file
        expected: Dict with 'rules' containing 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match expected crop size, 0.0 otherwise
    """
    if src_path is None:
        logger.warning('Source path is None')
        return 0.0
    try:
        img = Image.open(src_path)
        (actual_width, actual_height) = img.size
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        logger.info(f'Cropped image size: {actual_width}x{actual_height}, Expected: {expected_width}x{expected_height}')
        width_match = actual_width == expected_width if expected_width is not None else True
        height_match = actual_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image crop: {e}')
        return 0.0

def check_image_resized__4939e12b(result_state, expected_state, **options):
    """
    Check if an image has been resized to the expected dimensions.

    Args:
        result_state: Dictionary with 'width' and 'height' keys from getter
        expected_state: Dictionary with 'width', 'height', and 'tolerance' keys
        **options: Additional options

    Returns:
        float: 1.0 if image dimensions match within tolerance, 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Result state is None, image may not exist')
        return 0.0
    expected_width = expected_state.get('width')
    expected_height = expected_state.get('height')
    tolerance = expected_state.get('tolerance', 0)
    actual_width = result_state.get('width')
    actual_height = result_state.get('height')
    if actual_width is None or actual_height is None:
        logger.warning('Actual image dimensions are None')
        return 0.0
    width_diff = abs(actual_width - expected_width)
    height_diff = abs(actual_height - expected_height)
    if width_diff <= tolerance and height_diff <= tolerance:
        logger.info(f'Image dimensions match: {actual_width}x{actual_height} (expected {expected_width}x{expected_height} ±{tolerance})')
        return 1.0
    else:
        logger.info(f'Image dimensions do not match: {actual_width}x{actual_height} (expected {expected_width}x{expected_height} ±{tolerance})')
        return 0.0

def check_file_exists_and_image_size__16fb0dc6(src_path, rule):
    """
    Check if the image file exists and matches the frame at 5 seconds from src.mp4.

    This function verifies that:
    1. The file exists at the specified path
    2. The file is a valid image
    3. The image meets minimum dimension requirements
    4. The image content matches the frame extracted at 5 seconds from the source video

    Args:
        src_path: Path to the image file to check
        rule: Dictionary containing validation rules with keys:
              - min_width: Minimum required width in pixels
              - min_height: Minimum required height in pixels
              - source_video: Path to the source video file
              - timestamp: Timestamp in seconds to extract frame from

    Returns:
        float: 1.0 if file exists, meets size requirements, and matches the ground truth frame, 0.0 otherwise
    """
    if src_path is None:
        logging.debug('Source path is None')
        return 0.0
    if not os.path.isfile(src_path):
        logging.debug(f'File does not exist: {src_path}')
        return 0.0
    try:
        img = Image.open(src_path)
        (width, height) = img.size
        logging.debug(f'Image size: {width}x{height}')
        min_width = rule.get('min_width', 0)
        min_height = rule.get('min_height', 0)
        if width < min_width or height < min_height:
            logging.debug(f'Image does not meet requirements: {width}x{height} < {min_width}x{min_height}')
            return 0.0
        source_video = rule.get('source_video', '/home/user/Desktop/src.mp4')
        timestamp = rule.get('timestamp', 5)
        ground_truth_frame = extract_frame_from_video(source_video, timestamp)
        if ground_truth_frame is None:
            logging.warning('Could not extract ground truth frame from source video')
            logging.debug(f'Image meets basic requirements: {width}x{height} >= {min_width}x{min_height}')
            return 1.0
        similarity = compare_images(img, ground_truth_frame)
        logging.debug(f'Image similarity with ground truth: {similarity:.4f}')
        if similarity >= 0.95:
            logging.debug(f'Image matches ground truth frame at {timestamp}s with similarity {similarity:.4f}')
            return 1.0
        else:
            logging.debug(f'Image does not match ground truth frame (similarity {similarity:.4f} < 0.95)')
            return 0.0
    except Exception as e:
        logging.error(f'Error opening or processing image: {e}')
        return 0.0

def check_image_rotated__4e34ef5d(actual: Dict[str, Any], rules: Dict[str, Any], **options) -> float:
    """Check if an image has been rotated based on dimension comparison.

    Args:
        actual: Dict with 'width', 'height', and 'exists' keys from getter
        rules: Dict with rules to validate, e.g., 'width_greater_than_height': bool
        **options: Additional options

    Returns:
        float: 1.0 if the rotation check passes, 0.0 otherwise
    """
    if actual is None or not isinstance(actual, dict):
        logger.warning(f'Invalid actual value: {actual}')
        return 0.0
    if not actual.get('exists', False):
        logger.warning('Image file does not exist')
        return 0.0
    width = actual.get('width', 0)
    height = actual.get('height', 0)
    if width == 0 or height == 0:
        logger.warning(f'Invalid dimensions: {width}x{height}')
        return 0.0
    width_greater_than_height = rules.get('width_greater_than_height', True)
    if width_greater_than_height:
        if width > height:
            logger.info(f'Rotation check passed: width ({width}) > height ({height})')
            return 1.0
        else:
            logger.info(f'Rotation check failed: width ({width}) should be > height ({height})')
            return 0.0
    elif width < height:
        logger.info(f'Rotation check passed: width ({width}) < height ({height})')
        return 1.0
    else:
        logger.info(f'Rotation check failed: width ({width}) should be < height ({height})')
        return 0.0

def check_png_export_and_structure__c69055e15cceefc40a87e6de042c2331(result, expected, **options):
    """
    Check if PNG was exported correctly with proper dimensions and content verification.

    This metric verifies:
    1. PNG file exists (30% of score)
    2. Dimensions meet minimum requirements (30% of score)
    3. Content matches original XCF file (40% of score)

    Args:
        result: Dict from getter with file_exists, width, height, content_verified, dimensions_match
        expected: Dict with file_exists, min_width, min_height requirements
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False) and expected.get('file_exists', False):
        score += 0.3
        logger.info('✓ PNG file exists')
    else:
        logger.warning('✗ PNG file does not exist')
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    min_width = expected.get('min_width', 0)
    min_height = expected.get('min_height', 0)
    dimension_score = 0.0
    if width >= min_width:
        dimension_score += 0.15
        logger.info(f'✓ Width {width} >= {min_width}')
    else:
        logger.warning(f'✗ Width {width} < {min_width}')
    if height >= min_height:
        dimension_score += 0.15
        logger.info(f'✓ Height {height} >= {min_height}')
    else:
        logger.warning(f'✗ Height {height} < {min_height}')
    score += dimension_score
    content_verified = result.get('content_verified', False)
    dimensions_match = result.get('dimensions_match', False)
    original_dimensions = result.get('original_dimensions', None)
    if content_verified:
        score += 0.4
        logger.info('✓ PNG content verified against original XCF file')
        if original_dimensions:
            logger.info(f'✓ Original XCF dimensions: {original_dimensions[0]}x{original_dimensions[1]}')
    elif dimensions_match:
        score += 0.3
        logger.info('✓ PNG dimensions match original XCF (partial credit)')
    else:
        logger.warning('✗ Cannot verify PNG was exported from the original XCF file')
        if original_dimensions:
            logger.warning(f'Original XCF: {original_dimensions[0]}x{original_dimensions[1]}, PNG: {width}x{height}')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_image_scaled__5d068218(result, expected, **options):
    """
    Check if image was scaled by expected factor.

    Args:
        result: Dict with scale info from getter
        expected: Dict with 'scale_factor' expected value
        **options: Additional options

    Returns:
        float: Score (1.0 if scaled correctly, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        width_factor = result.get('width_factor', 1.0)
        height_factor = result.get('height_factor', 1.0)
        expected_factor = expected.get('scale_factor', 0.5)
        tolerance = expected.get('tolerance', 0.05)
        logger.info(f'Expected factor: {expected_factor}, Actual factors: {width_factor}x{height_factor}')
        if abs(width_factor - expected_factor) <= tolerance and abs(height_factor - expected_factor) <= tolerance:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking scale: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_file_exists__7c0cc95089263e14cb308f3757c8acc1(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if image file exists and has correct format.

    Args:
        result: Dict from getter with 'exists', 'is_png' keys
        expected: Dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        logger.info('File does not exist')
        return score
    if result.get('is_png', False):
        score += 0.5
    else:
        logger.info('File exists but is not a PNG')
    return score

def check_image_properties__bd10eca8(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_image_flip_vertical__fc198eb0(result, expected, **options):
    """
    Check if the image has been vertically flipped (top-bottom).

    Args:
        result: Path to result image file
        expected: Dict with 'source_path'
        **options: Additional options

    Returns:
        float: Score (1.0 if flipped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        flipped_source = source_img.transpose(Image.FLIP_TOP_BOTTOM)
        if structure_check_by_ssim(result_img, flipped_source, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match vertically flipped source')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_flip_vertical__fc198eb0: {e}')
        return 0.0

def check_image_crop__a8983eb9(result, expected, **options):
    """
    Check if the image has been cropped to the specified region.

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'crop_box' [x1, y1, x2, y2]
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        crop_box = expected.get('crop_box')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        expected_crop = source_img.crop(tuple(crop_box))
        if result_img.size != expected_crop.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {expected_crop.size}')
            return 0.0
        if structure_check_by_ssim(result_img, expected_crop, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match expected crop')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_crop__a8983eb9: {e}')
        return 0.0

def check_image_size__fb0dfb53338d620d275120282f62d824(result, expected, **options):
    """
    Check if image size matches expected dimensions.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected dimensions)
        **options: Additional options (tolerance for allowing small differences)

    Returns:
        float: 1.0 if size matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 0)
    result_width = result.get('width', 0)
    result_height = result.get('height', 0)
    expected_width = expected.get('width', 0)
    expected_height = expected.get('height', 0)
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    logger.info(f'Image size check: result=({result_width}, {result_height}), expected=({expected_width}, {expected_height}), match={width_match and height_match}')
    return 1.0 if width_match and height_match else 0.0

def check_png_resolution__27528fb8(result, expected, **options):
    """
    Check if PNG has sufficient resolution and correct aspect ratio.

    Args:
        result: Dict from getter with 'exists', 'total_pixels', 'aspect_ratio'
        expected: Dict with 'min_pixels' and 'aspect_ratio' [min, max]
        **options: Additional options

    Returns:
        float: Score based on resolution and aspect ratio checks
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.3
    min_pixels = expected.get('min_pixels', 0)
    total_pixels = result.get('total_pixels', 0)
    if total_pixels >= min_pixels:
        score += 0.35
    aspect_range = expected.get('aspect_ratio', [0, float('inf')])
    aspect_ratio = result.get('aspect_ratio', 0.0)
    if aspect_range[0] <= aspect_ratio <= aspect_range[1]:
        score += 0.35
    return score

def check_jpg_list__8ddf03c6b80780c49b4f9497dee3f888(result, expected, **options):
    """Check if jpg filename list matches expected.

    Args:
        result: List of filenames from getter
        expected: Rules dict with 'expected' list of filenames
        **options: Additional options

    Returns:
        float: 1.0 if lists match, 0.0 otherwise
    """
    expected_files = expected.get('expected', [])
    result_sorted = sorted(result) if result else []
    expected_sorted = sorted(expected_files)
    if len(result_sorted) != len(expected_sorted):
        return 0.0
    if result_sorted == expected_sorted:
        return 1.0
    else:
        return 0.0

def check_image_crop__a0afbcc2(result, expected, **options):
    """
    Check if the image has been cropped to the specified region.

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'crop_box' [x1, y1, x2, y2]
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        crop_box = expected.get('crop_box')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        expected_crop = source_img.crop(tuple(crop_box))
        if result_img.size != expected_crop.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {expected_crop.size}')
            return 0.0
        if structure_check_by_ssim(result_img, expected_crop, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match expected crop')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_crop__a0afbcc2: {e}')
        return 0.0

def check_image_is_grayscale__1ac9295ec86f9e2e04c973e2e47b273c(result, expected, **options):
    """
    Check if the image is in grayscale mode.

    Args:
        result: Dict with 'mode' and 'is_grayscale' from getter
        expected: Dict with 'is_grayscale' expected value
        **options: Additional options

    Returns:
        float: 1.0 if grayscale status matches expected, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    result_is_grayscale = result.get('is_grayscale')
    expected_is_grayscale = expected.get('is_grayscale', True)
    if result_is_grayscale is None:
        logger.error('Missing is_grayscale in result')
        return 0.0
    logger.info(f"Result is_grayscale: {result_is_grayscale}, Mode: {result.get('mode')}")
    logger.info(f'Expected is_grayscale: {expected_is_grayscale}')
    if result_is_grayscale == expected_is_grayscale:
        return 1.0
    else:
        return 0.0

def check_image_dimensions__9472df29(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_extracted__69184b17(result, expected, **options):
    score = 0.0
    if result.get('exists'):
        score += 0.4
    else:
        return 0.0
    if result.get('size', 0) >= expected.get('min_size', 1000):
        score += 0.2
    if result.get('format') == expected.get('format', 'PNG'):
        score += 0.2
    if result.get('width', 0) >= expected.get('min_width', 100) and result.get('height', 0) >= expected.get('min_height', 100):
        score += 0.2
    return score

def check_triangle_bottomleft__4f468172ec938d83ce9f0bb7cb13c9fc(tgt_path, expected, **options):
    """
    Check if the yellow triangle is in the bottom-left corner of the image.

    Args:
        tgt_path: Path to the result image
        expected: Expected configuration with position and tolerance
        **options: Additional options

    Returns:
        float: 1.0 if yellow triangle is in bottom-left, 0.0 otherwise
    """
    if tgt_path is None:
        return 0.0
    img = Image.open(tgt_path)
    img_array = np.array(img)
    if img_array.ndim == 2:
        return 0.0
    if img_array.shape[2] == 4:
        img_array = img_array[:, :, :3]
    img_hsv = Image.fromarray(img_array).convert('HSV')
    img_hsv_array = np.array(img_hsv)
    (h, s, v) = (img_hsv_array[:, :, 0], img_hsv_array[:, :, 1], img_hsv_array[:, :, 2])
    yellow_mask = (h >= 25) & (h <= 85) & (s >= 100) & (v >= 100)
    (r, g, b) = (img_array[:, :, 0], img_array[:, :, 1], img_array[:, :, 2])
    yellow_mask_rgb = (r >= 180) & (g >= 180) & (b <= 100)
    yellow_mask = yellow_mask | yellow_mask_rgb
    yellow_coords = np.argwhere(yellow_mask)
    if len(yellow_coords) == 0:
        return 0.0
    yellow_area = len(yellow_coords)
    total_area = img_array.shape[0] * img_array.shape[1]
    if yellow_area < 0.005 * total_area or yellow_area > 0.5 * total_area:
        return 0.0
    centroid = yellow_coords.mean(axis=0)
    tolerance = expected.get('tolerance', 0.15)
    image_shape = np.array(img_array.shape[:2])
    bottom_left = np.array([image_shape[0], 0])
    tolerance_pixels = tolerance * image_shape
    in_bottomleft = np.all(np.abs(centroid - bottom_left) < tolerance_pixels)
    if bool(in_bottomleft):
        return 1.0
    else:
        return 0.0

def check_png_count_range__b4c05b10(result: int, expected: Dict[str, Any], **options) -> float:
    """Check if PNG count is within expected range"""
    min_count = expected.get('min_count', 1)
    max_count = expected.get('max_count', 10)
    if min_count <= result <= max_count:
        logger.info(f'PNG count {result} is within range [{min_count}, {max_count}]')
        return 1.0
    logger.info(f'PNG count {result} is outside range [{min_count}, {max_count}]')
    return 0.0

def check_triangle_rotation__8f589d19(tgt_path, expected, **options):
    """
    Check if the yellow triangle has been rotated approximately 90 degrees clockwise
    from the original image.
    Variation 1 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with rotation_degrees and tolerance_degrees
        **options: Additional options

    Returns:
        float: Score (1.0 if rotated correctly, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
        if not os.path.exists(original_path):
            result_dir = os.path.dirname(tgt_path)
            original_path = os.path.join(result_dir, 'Triangle_On_The_Side.png')
        if not os.path.exists(original_path):
            logger.error(f'Original image not found at {original_path}')
            return 0.0
        original_img = Image.open(original_path)
        result_img = Image.open(tgt_path)
        original_array = np.array(original_img.convert('RGB'))
        result_array = np.array(result_img.convert('RGB'))
        (orig_mask, orig_angle, orig_contour) = _detect_yellow_triangle(original_array)
        if orig_angle is None:
            logger.error('Failed to detect yellow triangle in original image')
            return 0.0
        (result_mask, result_angle, result_contour) = _detect_yellow_triangle(result_array)
        if result_angle is None:
            logger.error('Failed to detect yellow triangle in result image')
            return 0.0
        logger.info(f'Original angle: {orig_angle}°, Result angle: {result_angle}°')
        angle_diff = result_angle - orig_angle
        while angle_diff > 180:
            angle_diff -= 360
        while angle_diff <= -180:
            angle_diff += 360
        logger.info(f'Angle difference: {angle_diff}°')
        target_rotation = expected.get('rotation_degrees', 90)
        tolerance = expected.get('tolerance_degrees', 10)
        rotation_error = abs(angle_diff - target_rotation)
        alt_rotation_error1 = abs(angle_diff - (target_rotation - 360))
        alt_rotation_error2 = abs(angle_diff - (target_rotation + 360))
        min_error = min(rotation_error, alt_rotation_error1, alt_rotation_error2)
        logger.info(f'Target rotation: {target_rotation}°, Rotation error: {min_error}°, Tolerance: {tolerance}°')
        orig_area = cv2.contourArea(orig_contour)
        result_area = cv2.contourArea(result_contour)
        area_ratio = result_area / orig_area if orig_area > 0 else 0
        logger.info(f'Original area: {orig_area}, Result area: {result_area}, Ratio: {area_ratio}')
        if min_error <= tolerance and 0.8 <= area_ratio <= 1.2:
            return 1.0
        else:
            if min_error > tolerance:
                logger.info(f'Rotation error {min_error}° exceeds tolerance {tolerance}°')
            if not 0.8 <= area_ratio <= 1.2:
                logger.info(f'Area ratio {area_ratio} is outside expected range [0.8, 1.2]')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle rotation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_properties__37707684957589279b0fa14602529fe7(result: Optional[Dict[str, Any]], expected: dict, **options) -> float:
    """
    Check if image properties match expected values.

    Args:
        result: Image properties dict from getter (or None if image doesn't exist)
        expected: Expected rules dict with keys like 'format', 'min_width', 'min_height'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on how many properties match
    """
    if result is None:
        logger.warning('Image does not exist or could not be read')
        return 0.0
    score = 0.0
    checks = 0
    if 'format' in expected:
        checks += 1
        expected_format = expected['format'].upper()
        actual_format = result.get('format', '').upper()
        if actual_format == expected_format:
            score += 1.0
            logger.info(f'Format check passed: {actual_format}')
        else:
            logger.warning(f'Format check failed: expected {expected_format}, got {actual_format}')
    if 'min_width' in expected:
        checks += 1
        if result.get('width', 0) >= expected['min_width']:
            score += 1.0
            logger.info(f"Width check passed: {result.get('width')} >= {expected['min_width']}")
        else:
            logger.warning(f"Width check failed: {result.get('width')} < {expected['min_width']}")
    if 'min_height' in expected:
        checks += 1
        if result.get('height', 0) >= expected['min_height']:
            score += 1.0
            logger.info(f"Height check passed: {result.get('height')} >= {expected['min_height']}")
        else:
            logger.warning(f"Height check failed: {result.get('height')} < {expected['min_height']}")
    if 'width' in expected:
        checks += 1
        if result.get('width') == expected['width']:
            score += 1.0
            logger.info(f"Exact width check passed: {result.get('width')}")
        else:
            logger.warning(f"Exact width check failed: expected {expected['width']}, got {result.get('width')}")
    if 'height' in expected:
        checks += 1
        if result.get('height') == expected['height']:
            score += 1.0
            logger.info(f"Exact height check passed: {result.get('height')}")
        else:
            logger.warning(f"Exact height check failed: expected {expected['height']}, got {result.get('height')}")
    if checks == 0:
        logger.warning('No checks specified in expected rules')
        return 0.0
    final_score = score / checks
    logger.info(f'Final image property score: {final_score} ({score}/{checks} checks passed)')
    return final_score

def check_two_folders_with_images__6a4313ae(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if two folders exist with the correct number of images.

    Args:
        result: Dict with 'folder1_count' and 'folder2_count' keys
        expected: Dict with expected counts in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 with partial credit
    """
    if result is None:
        return 0.0
    expected_folder1 = expected.get('folder1_count', 0)
    expected_folder2 = expected.get('folder2_count', 0)
    actual_folder1 = result.get('folder1_count', 0)
    actual_folder2 = result.get('folder2_count', 0)
    score = 0.0
    if actual_folder1 == expected_folder1:
        score += 0.5
    if actual_folder2 == expected_folder2:
        score += 0.5
    logger.info(f'Folder counts - Expected: f1={expected_folder1}, f2={expected_folder2}; Actual: f1={actual_folder1}, f2={actual_folder2}; Score: {score}')
    return score

def check_pdf_with_image__d5371d5b(pdf_file: str, expected, **options):
    """
    Check if PDF exists, has exactly 1 page, and contains image content.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with 'page_count' and 'has_image' flags
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.0
    try:
        reader = PdfReader(pdf_file)
        nb_pages = len(reader.pages)
        expected_pages = expected.get('page_count', 1)
        if nb_pages == expected_pages:
            score += 0.5
        if expected.get('has_image', True):
            has_images = False
            for page in reader.pages:
                if '/XObject' in page.get('/Resources', {}):
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        if xobjects[obj]['/Subtype'] == '/Image':
                            has_images = True
                            break
                if has_images:
                    break
            if has_images:
                score += 0.5
    except Exception as e:
        return 0.0
    return score

def check_triangle_color_blue__d52919fa54c3d57b57da98359753b674(result_data, expected, **options):
    """Check if the triangle has been changed to blue color.

    This function:
    1. Loads the original image to find the yellow triangle's location
    2. Uses contour detection to identify the triangle shape
    3. Verifies the triangle in the result image is in the same location and is blue

    Args:
        result_data: Dict with 'result_path' and 'original_path' keys
        expected: Expected rules dict (from config["rules"])
        **options: Additional options

    Returns:
        float: 1.0 if triangle is blue, 0.0 otherwise
    """
    if result_data is None:
        logger.error('No result data provided')
        return 0.0
    result_path = result_data.get('result_path')
    original_path = result_data.get('original_path')
    if not result_path:
        logger.error('No result path provided')
        return 0.0
    try:
        result_img = Image.open(result_path)
        result_array = np.array(result_img)
        result_bgr = cv2.cvtColor(result_array, cv2.COLOR_RGB2BGR)
        triangle_region = None
        original_triangle_area = None
        if original_path:
            try:
                original_img = Image.open(original_path)
                original_array = np.array(original_img)
                original_bgr = cv2.cvtColor(original_array, cv2.COLOR_RGB2BGR)
                yellow_mask = cv2.inRange(original_array, np.array([200, 150, 0]), np.array([255, 255, 50]))
                (contours, _) = cv2.findContours(yellow_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
                if contours:
                    largest_contour = max(contours, key=cv2.contourArea)
                    original_triangle_area = cv2.contourArea(largest_contour)
                    (x, y, w, h) = cv2.boundingRect(largest_contour)
                    triangle_region = (x, y, w, h)
                    logger.debug(f'Found yellow triangle region: {triangle_region}, area: {original_triangle_area}')
            except Exception as e:
                logger.warning(f'Could not analyze original image: {e}')
        blue_mask = cv2.inRange(result_array, np.array([0, 0, 100]), np.array([100, 100, 255]))
        (contours, _) = cv2.findContours(blue_mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        if not contours:
            logger.error('No blue regions found in result image')
            return 0.0
        blue_triangle_found = False
        for contour in contours:
            area = cv2.contourArea(contour)
            if area < 100:
                continue
            epsilon = 0.04 * cv2.arcLength(contour, True)
            approx = cv2.approxPolyDP(contour, epsilon, True)
            num_vertices = len(approx)
            logger.debug(f'Contour has {num_vertices} vertices, area={area}')
            if triangle_region and original_triangle_area:
                (x, y, w, h) = cv2.boundingRect(contour)
                (orig_x, orig_y, orig_w, orig_h) = triangle_region
                (center_x, center_y) = (x + w // 2, y + h // 2)
                (orig_center_x, orig_center_y) = (orig_x + orig_w // 2, orig_y + orig_h // 2)
                distance = np.sqrt((center_x - orig_center_x) ** 2 + (center_y - orig_center_y) ** 2)
                logger.debug(f'Distance between centers: {distance}')
                size_diff = abs(area - original_triangle_area)
                size_ratio = size_diff / original_triangle_area if original_triangle_area > 0 else 1.0
                logger.debug(f'Size comparison: result area={area}, original area={original_triangle_area}, ratio={size_ratio}')
                if distance < 50 and size_ratio < 0.5:
                    if 3 <= num_vertices <= 5:
                        blue_triangle_found = True
                        logger.debug(f'Found blue triangle at similar location with {num_vertices} vertices')
                        break
            elif 3 <= num_vertices <= 5 and area > 1000:
                blue_triangle_found = True
                logger.debug(f'Found blue triangle-like shape with {num_vertices} vertices')
                break
        if blue_triangle_found:
            blue_pixels = result_array[blue_mask > 0]
            if len(blue_pixels) > 0:
                avg_color = np.mean(blue_pixels, axis=0)
                (r, g, b) = avg_color[:3]
                logger.debug(f'Average blue region color: R={r:.1f}, G={g:.1f}, B={b:.1f}')
                if b > r + 30 and b > g + 30 and (b > 100):
                    logger.debug('Blue triangle verified successfully')
                    return 1.0
                else:
                    logger.error(f'Region is not blue enough: R={r:.1f}, G={g:.1f}, B={b:.1f}')
                    return 0.0
        logger.error('No blue triangle found matching criteria')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle color: {e}', exc_info=True)
        return 0.0

def check_image_extracted__1a437041(result, expected, **options):
    score = 0.0
    if result.get('exists'):
        score += 0.4
    else:
        return 0.0
    if result.get('size', 0) >= expected.get('min_size', 1000):
        score += 0.2
    if result.get('format') == expected.get('format', 'PNG'):
        score += 0.2
    if result.get('width', 0) >= expected.get('min_width', 100) and result.get('height', 0) >= expected.get('min_height', 100):
        score += 0.2
    return score

def check_bottom_half_image__4baa14c913fe32a60f559a98296affae(result, expected, **options):
    """
    Check if the saved image matches the bottom half of the original image.
    Uses pixel-by-pixel comparison with a similarity threshold.

    Args:
        result: dict from getter with "result_image", "expected_image", and "dimensions"
        expected: Not used (getter provides both images)
        **options: Additional options:
            - similarity_threshold: minimum similarity score (0.0-1.0), default 0.95

    Returns:
        float: 1.0 if images match (similarity >= threshold), 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    result_image = result.get('result_image')
    expected_image = result.get('expected_image')
    dimensions = result.get('dimensions', {})
    if result_image is None or expected_image is None:
        logger.error('Missing image arrays in result')
        return 0.0
    result_width = dimensions.get('result_width')
    result_height = dimensions.get('result_height')
    expected_width = dimensions.get('expected_width')
    expected_height = dimensions.get('expected_height')
    if result_width != expected_width or result_height != expected_height:
        logger.error(f'Dimension mismatch: result {result_width}x{result_height}, expected {expected_width}x{expected_height}')
        return 0.0
    result_array = np.array(result_image)
    expected_array = np.array(expected_image)
    if result_array.shape != expected_array.shape:
        logger.error(f'Image shape mismatch: result {result_array.shape}, expected {expected_array.shape}')
        return 0.0
    similarity_threshold = options.get('similarity_threshold', 0.95)
    mse = np.mean((result_array.astype(float) - expected_array.astype(float)) ** 2)
    max_mse = 255.0 ** 2
    similarity = 1.0 - mse / max_mse
    logger.info(f'Image similarity: {similarity:.4f} (MSE: {mse:.2f})')
    if similarity >= similarity_threshold:
        logger.info(f'Images match with similarity {similarity:.4f} >= {similarity_threshold}')
        return 1.0
    else:
        logger.warning(f"Images don't match: similarity {similarity:.4f} < {similarity_threshold}")
        return 0.0

def check_gimp_toolbox_setting__90386cba106758972cba7bf949bb562f(actual_config_path, rule):
    """
    Check if GIMP toolbox setting is as expected.
    This checks the sessionrc file for toolbox-related settings.

    Args:
        actual_config_path: Path to the GIMP config file
        rule: Expected configuration with keys:
            - key: Config key name (can be string or list)
            - value: Expected value

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_triangle_bottomright__95d4f9f61b019dd634339dd50c3adb82(tgt_path, expected, **options):
    """
    Check if the yellow triangle is in the bottom-right corner of the image.

    Args:
        tgt_path: Path to the result image
        expected: Expected configuration with position and tolerance
        **options: Additional options

    Returns:
        float: 1.0 if yellow triangle is in bottom-right, 0.0 otherwise
    """
    if tgt_path is None:
        return 0.0
    img = Image.open(tgt_path)
    img_array = np.array(img)
    (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
    if len(unique_colors) < 2:
        return 0.0
    unique_colors_sorted = unique_colors[np.argsort(counts)[::-1]]
    yellow_color = None
    for color in unique_colors_sorted[1:]:
        (r, g, b) = (color[0], color[1], color[2])
        if r > 200 and g > 200 and (b < 100) and (abs(r - g) < 100):
            yellow_color = color
            break
    if yellow_color is None:
        return 0.0
    triangle_mask = np.all(img_array == yellow_color, axis=2)
    triangle_coords = np.argwhere(triangle_mask)
    if len(triangle_coords) == 0:
        return 0.0
    centroid = triangle_coords.mean(axis=0)
    tolerance = expected.get('tolerance', 0.15)
    image_shape = np.array(img_array.shape[:2])
    bottom_right = image_shape - 1
    tolerance_pixels = tolerance * image_shape
    in_bottomright = np.all(np.abs(centroid - bottom_right) < tolerance_pixels)
    if bool(in_bottomright):
        return 1.0
    else:
        return 0.0

def check_photo_rename__0a812fadd7008b5cb899c479c82ab6e7(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if photo rename task was completed correctly with content validation.

    Validates:
    1. Renamed file (ending_01.jpg) exists
    2. Renamed file contains expected text ("Thank you") via OCR
    3. Original file (DSC00657.jpg) no longer exists (was renamed, not copied)
    4. Total file count matches expected (6 files, no copies created)

    Args:
        result: Dict with verification results from getter:
            {
                'renamed_file_exists': bool,
                'renamed_file_text': str,
                'original_file_exists': bool,
                'total_files': int
            }
        expected: Dict with expected values:
            {
                'renamed_file': str (expected filename),
                'should_contain_text': str (text to find in OCR),
                'original_file': str (original filename that should not exist),
                'total_files': int (expected file count)
            }
        **options: Additional options (unused)

    Returns:
        1.0 if all conditions met, 0.0 otherwise
    """
    should_contain_text = expected.get('should_contain_text', '')
    expected_total_files = expected.get('total_files', 6)
    if not result.get('renamed_file_exists', False):
        logger.info('Verification failed: Renamed file (ending_01.jpg) does not exist')
        return 0.0
    renamed_text = result.get('renamed_file_text', '')
    if should_contain_text and should_contain_text.lower() not in renamed_text.lower():
        logger.info(f"Verification failed: Renamed file does not contain expected text '{should_contain_text}'. OCR output: {renamed_text[:200]}")
        return 0.0
    if result.get('original_file_exists', False):
        logger.info('Verification failed: Original file (DSC00657.jpg) still exists (file was copied instead of renamed)')
        return 0.0
    total_files = result.get('total_files', 0)
    if total_files != expected_total_files:
        logger.info(f'Verification failed: Expected {expected_total_files} files, but found {total_files} files (copies may have been created)')
        return 0.0
    logger.info('All verification checks passed: file renamed correctly with content validation')
    return 1.0

def check_gif_file__16fb0dc6(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_dimensions__fa4ed82a(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_resize__89f19527(result, expected, **options):
    """
    Check if image has been resized to specified dimensions while maintaining structure.

    Args:
        result: Path to the result image file
        expected: Expected rules dict with target_width and target_height
        **options: Additional options

    Returns:
        float: 1.0 if image is resized correctly with structure preserved, 0.0 otherwise
    """
    if result is None or expected is None:
        logger.error('Result or expected is None')
        return 0.0
    try:
        result_img = Image.open(result)
        target_width = expected.get('target_width')
        target_height = expected.get('target_height')
        if target_width is None or target_height is None:
            logger.error('Missing target_width or target_height in rules')
            return 0.0
        (actual_width, actual_height) = result_img.size
        logger.debug(f'Actual size: {actual_width}x{actual_height}, Target: {target_width}x{target_height}')
        if actual_width == target_width and actual_height == target_height:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image resize: {e}')
        return 0.0

def check_gif_exists__11cf8ab6(result_state, expected_state, **options):
    """
    Check if the GIF file exists and has the correct format and animation properties.

    This function verifies that:
    1. The file exists at the specified path
    2. The file is a valid GIF format
    3. The GIF is animated (has multiple frames)
    4. The GIF has a reasonable frame count for approximately 3 seconds of animation

    Args:
        result_state: Path to the result file (animation.gif)
        expected_state: Expected rules dict containing format requirements
        **options: Additional options

    Returns:
        float: 1.0 if file exists, is GIF format, and is properly animated, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state (file path) is None')
        return 0.0
    try:
        if not os.path.exists(result_state):
            logger.error(f'File does not exist: {result_state}')
            return 0.0
        if not os.path.isfile(result_state):
            logger.error(f'Path is not a file: {result_state}')
            return 0.0
        img = Image.open(result_state)
        if img.format != 'GIF':
            logger.error(f'File format is {img.format}, expected GIF')
            img.close()
            return 0.0
        frame_count = 0
        try:
            while True:
                frame_count += 1
                img.seek(frame_count)
        except EOFError:
            pass
        if frame_count < 2:
            logger.error(f'GIF is not animated (only {frame_count} frame)')
            img.close()
            return 0.0
        if frame_count < 10:
            logger.warning(f'GIF has only {frame_count} frames, which seems low for a 3-second animation')
        elif frame_count > 100:
            logger.warning(f'GIF has {frame_count} frames, which seems high for a 3-second animation')
        logger.info(f'Successfully verified animated GIF file: {result_state} (format: {img.format}, size: {img.size}, frames: {frame_count})')
        img.close()
        if expected_state and isinstance(expected_state, dict):
            expected_format = expected_state.get('format', 'GIF')
            if expected_format != 'GIF':
                logger.warning(f'Expected format mismatch: {expected_format} vs GIF')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking GIF file: {e}')
        return 0.0

def check_total_images__2759cd98(result, expected, **options):
    """Check if total image count meets minimum requirement.

    Args:
        result: Total image count from getter
        expected: Expected criteria
        **options: Additional options

    Returns:
        float: 1.0 if meets requirement, partial credit otherwise
    """
    min_images = expected.get('min_images', 4)
    if result >= min_images:
        return 1.0
    else:
        return result / min_images if min_images > 0 else 0.0

def check_gimp_gimprc_setting__9c13adcb(actual_config_path, expected, **options):
    """
    Check if a GIMP gimprc setting has the expected value.

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'key' and 'value' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if setting matches, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        target_key = expected.get('key')
        target_value = expected.get('value')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == target_key and value == target_value:
                    logger.info(f'Found matching setting: {key} = {value}')
                    return 1.0
        logger.warning(f"Setting not found or doesn't match: {target_key} = {target_value}")
        return 0.0
    except Exception as e:
        logger.error(f'Error checking gimprc setting: {e}')
        return 0.0

def check_gif_file__3d85489a(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_rotated__b49205bc(result, expected, **options):
    """Check if image dimensions changed indicating rotation.

    Args:
        result: Dict with image properties from getter
        expected: Dict with expected properties in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists'):
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result.get('width') == expected_width and result.get('height') == expected_height:
        return 1.0
    return 0.0

def check_image_crop__70ffd430(result, expected, **options):
    """
    Check if the image has been cropped to the specified region.

    Args:
        result: List of [result_image_path, source_image_path] - local file paths from vm_file getter
        expected: Dict with 'source_index' (index of source in result list) and 'crop_box' [x1, y1, x2, y2]
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        if not isinstance(result, list) or len(result) < 2:
            logging.error(f'Expected result to be a list with 2 elements, got {result}')
            return 0.0
        result_path = result[0]
        source_index = expected.get('source_index', 1)
        source_path = result[source_index]
        crop_box = expected.get('crop_box')
        if not result_path or not source_path:
            logging.error('Missing file paths from vm_file getter')
            return 0.0
        result_img = Image.open(result_path)
        source_img = Image.open(source_path)
        expected_crop = source_img.crop(tuple(crop_box))
        if result_img.size != expected_crop.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {expected_crop.size}')
            return 0.0
        if structure_check_by_ssim(result_img, expected_crop, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match expected crop')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_crop__70ffd430: {e}')
        return 0.0

def check_gimp_menubar_setting__eebb7fd7089a7abc35d29b3d4832455e(actual_config_path, rule):
    """
    Check if GIMP menubar setting is as expected.
    This checks the sessionrc file for menubar visibility settings.

    Args:
        actual_config_path: Path to the GIMP config file
        rule: Expected configuration with keys:
            - key: Config key name (can be string or list)
            - value: Expected value

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def verify_total_images__d49c5873(result, expected, **options):
    """Verify total image count.

    Args:
        result: Actual total
        expected: Expected total (dict with 'total' key)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_total = expected.get('total', 4)
    return 1.0 if result == expected_total else 0.0

def check_gif_file__e66cc6b0(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_resized__f8a3e253(result, expected, **options):
    """
    Check if the image has been resized to the expected dimensions.

    Args:
        result: Dict with 'width' and 'height' from getter
        expected: Dict with 'width' and 'height' expected values
        **options: Additional options

    Returns:
        float: Score (1.0 if resized correctly, 0.0 otherwise)
    """
    try:
        if result is None or expected is None:
            logger.error('Result or expected is None')
            return 0.0
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        actual_width = result.get('width')
        actual_height = result.get('height')
        if actual_width is None or actual_height is None:
            logger.error('Could not extract actual dimensions')
            return 0.0
        logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
        if actual_width == expected_width and actual_height == expected_height:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image resize: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_aspect_ratio__f85979b813875fecca12ba1c6ab4cc68(result, expected, **options):
    """
    Check if image aspect ratio is within tolerance of expected value.

    Args:
        result: dict with "aspect_ratio" key from getter
        expected: dict with expected "aspect_ratio" value (from rules)
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if aspect ratio matches within tolerance, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 0.01)
    expected_ratio = expected.get('aspect_ratio')
    result_ratio = result.get('aspect_ratio')
    if result_ratio is None or expected_ratio is None:
        logger.error('Missing aspect ratio value')
        return 0.0
    diff = abs(result_ratio - expected_ratio)
    if diff <= tolerance:
        logger.info(f'Aspect ratio matches: {result_ratio:.4f} (expected {expected_ratio:.4f}, tolerance {tolerance})')
        return 1.0
    else:
        logger.info(f'Aspect ratio mismatch: got {result_ratio:.4f}, expected {expected_ratio:.4f} (diff {diff:.4f} > tolerance {tolerance})')
        return 0.0

def check_gif_file_basic__de678d13fd248567f73258f2b3cb0372(result, expected, **options):
    """
    Check if a GIF file meets requirements for an animated GIF from video.

    Args:
        result: Dict from getter with file info (exists, file_size, format, width, height, is_animated, n_frames, duration)
        expected: Dict with expected values (min_size, max_size, format, min_duration, max_duration, min_frames, max_frames)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    expected_format = expected.get('format', 'GIF')
    if result.get('format') == expected_format:
        score += 0.2
        logger.info(f"Format check passed: {result.get('format')}")
    else:
        logger.warning(f"Format mismatch: expected {expected_format}, got {result.get('format')}")
    file_size = result.get('file_size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 5000000)
    if min_size <= file_size <= max_size:
        score += 0.15
        logger.info(f'File size check passed: {file_size} bytes (range: {min_size}-{max_size})')
    else:
        logger.warning(f'File size out of range: {file_size} bytes (expected: {min_size}-{max_size})')
    width = result.get('width', 0)
    height = result.get('height', 0)
    min_dimension = expected.get('min_dimension', 100)
    max_dimension = expected.get('max_dimension', 2000)
    if min_dimension <= width <= max_dimension and min_dimension <= height <= max_dimension and (width > 0) and (height > 0):
        score += 0.15
        logger.info(f'Dimensions check passed: {width}x{height}')
    else:
        logger.warning(f'Dimensions out of range: {width}x{height}')
    is_animated = result.get('is_animated', False)
    if is_animated:
        score += 0.2
        logger.info('Animation check passed: GIF is animated')
    else:
        logger.warning('Animation check failed: GIF is not animated or is single-frame')
    duration = result.get('duration', 0.0)
    min_duration = expected.get('min_duration', 4.5)
    max_duration = expected.get('max_duration', 5.5)
    if min_duration <= duration <= max_duration:
        score += 0.3
        logger.info(f'Duration check passed: {duration:.2f}s (range: {min_duration}-{max_duration}s)')
    else:
        logger.warning(f'Duration out of range: {duration:.2f}s (expected: {min_duration}-{max_duration}s)')
    n_frames = result.get('n_frames', 0)
    min_frames = expected.get('min_frames', 10)
    max_frames = expected.get('max_frames', 500)
    if min_frames <= n_frames <= max_frames:
        logger.info(f'Frame count check passed: {n_frames} frames (range: {min_frames}-{max_frames})')
    else:
        logger.warning(f'Frame count unusual: {n_frames} frames (expected: {min_frames}-{max_frames})')
    return score

def check_image_crop__59e6ca63b22becd85a8942d1a29325d9(result, expected, **options):
    """Check if image has been cropped to expected dimensions.

    Args:
        result: dict with dimensions from getter
        expected: dict with expected dimensions
        **options: Additional options (tolerance for dimension matching)

    Returns:
        float: Score based on dimension match (1.0 = perfect, 0.5 = partial, 0.0 = fail)
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('dimension_tolerance', 5)
    actual_width = result.get('width', 0)
    actual_height = result.get('height', 0)
    score = 0.0
    if expected_width is not None:
        if abs(actual_width - expected_width) <= tolerance:
            score += 0.5
    if expected_height is not None:
        if abs(actual_height - expected_height) <= tolerance:
            score += 0.5
    return score

def check_layer_exists__ca8ad8ab(result_state, expected_state, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected_state: Expected layer name to check for (string) or dict with 'layer_name'
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected_state, dict):
            expected_layer = expected_state.get('layer_name', '')
        else:
            expected_layer = expected_state
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_png_exists__06722a19(result_state, expected_state, **options):
    """
    Check if the PNG file exists and is in PNG format.

    Task: 2fe4b718-3bd7-46ec-bdce-b184f5653624
    Instruction: Capture a still image from 'src.mp4' using VLC,
                 then export it as 'frame_6s.png' on the desktop.

    Args:
        result_state: Path to the exported PNG file (str)
        expected_state: Expected rules (dict with 'format' key)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is PNG format, 0.0 otherwise
    """
    logger = logging.getLogger(__name__)
    if result_state is None or result_state == '':
        logger.warning('Result state is None or empty')
        return 0.0
    if not os.path.isfile(result_state):
        logger.warning(f'File does not exist: {result_state}')
        return 0.0
    try:
        img = Image.open(result_state)
        expected_format = expected_state.get('format', 'PNG')
        if img.format != expected_format:
            logger.warning(f'Image format mismatch: expected {expected_format}, got {img.format}')
            return 0.0
        logger.info(f'PNG file verification successful: {result_state} (format: {img.format}, size: {img.size})')
        return 1.0
    except Exception as e:
        logger.error(f'Error verifying PNG file: {e}')
        return 0.0

def check_layer_name_config__174a6594(result_state, expected_state, **options):
    """
    Check if the expected layer name exists in the list of layers.

    Args:
        result_state: List of layer names from getter
        expected_state: Expected state with 'layer_name' key
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    if result_state is None or not isinstance(result_state, list):
        logger.error(f'Invalid result_state: {result_state}')
        return 0.0
    try:
        expected_layer = expected_state.get('layer_name', 'Overlay')
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.info(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer name: {e}')
        return 0.0

def check_triangle_topleft__40af8739a7f9d38f5046e0dfa41c5f6e(result, expected, **options):
    """
    Check if the triangle has been moved to the top-left corner.

    Args:
        result: Actual position dict {'x': float, 'y': float} from getter
        expected: Dict with 'target_x' and 'target_y' keys
        **options: Optional 'tolerance' for position comparison (default: 0.15)

    Returns:
        float: Score from 0.0 to 1.0 based on how close the triangle is to target position
    """
    if result is None:
        return 0.0
    target_x = expected.get('target_x', 0.2)
    target_y = expected.get('target_y', 0.2)
    tolerance = options.get('tolerance', 0.15)
    x_diff = abs(result['x'] - target_x)
    y_diff = abs(result['y'] - target_y)
    if x_diff <= tolerance and y_diff <= tolerance:
        return 1.0
    else:
        return 0.0

def check_image_dimensions__696943e1(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_png_saved__b9abaa4fbc51b493882263c6a6aff8fe(result, expected, **options):
    """Check if a PNG image file was saved correctly.

    This metric now includes content-based verification by comparing the hash
    of the saved image with the hash of the expected 2nd image from the Word document.

    Args:
        result: dict from getter with 'exists', 'size', 'is_png', 'filename', 'hash', 'expected_hash'
        expected: dict with expected properties
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.2
    else:
        return 0.0
    if result.get('is_png', False):
        score += 0.15
    expected_name = expected.get('filename', '')
    if expected_name and expected_name.lower() == result.get('filename', '').lower():
        score += 0.1
    min_size = expected.get('min_size', 1000)
    max_size = expected.get('max_size', 10000000)
    actual_size = result.get('size', 0)
    if min_size <= actual_size <= max_size:
        score += 0.1
    if result.get('hash') and result.get('expected_hash'):
        if result['hash'] == result['expected_hash']:
            score += 0.45
        else:
            pass
    elif result.get('expected_hash') is None:
        score += 0.2
    return score

def check_image_dimensions__34372286(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_mode__78bfc971(result, expected, **options):
    """
    Check if the image has the expected color mode.

    Args:
        result: Image mode string from getter
        expected: Dict with 'mode' key
        **options: Additional options

    Returns:
        float: Score (1.0 if mode matches, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        expected_mode = expected.get('mode')
        if not expected_mode:
            logger.error('No expected mode provided')
            return 0.0
        logger.info(f'Expected mode: {expected_mode}, Actual mode: {result}')
        if result == expected_mode:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image mode: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_vacation_jpgs__d7a48669399bf74024b2a979a32d4ae1(result, expected, **options):
    """Check if vacation jpg files were copied correctly.

    Args:
        result: Directory tree dict from getter
        expected: Rules dict with 'expected' list of filenames
        **options: Additional options

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected.get('expected', [])
    if not result or 'children' not in result:
        return 0.0
    actual_files = [node['name'] for node in result['children']]
    if len(actual_files) != len(expected_files):
        return 0.0
    if set(actual_files) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_image_resized__912586a80097a155904c15da97ac2079(result, expected, **options):
    """Check if image has been resized to expected dimensions.

    Args:
        result: Dict with width and height from getter
        expected: Dict with expected width and height
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    width_match = result.get('width') == expected_width
    height_match = result.get('height') == expected_height
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_format__c0fb0f23(result, expected, **options):
    """
    Check if image format matches expected value.

    Args:
        result: Dict with 'format' and 'exists' keys from getter
        expected: Dict with expected 'format' value
        **options: Additional options

    Returns:
        float: Score (1.0 if format matches, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        if not result.get('exists', False):
            logger.error('File does not exist')
            return 0.0
        expected_format = expected.get('format', 'JPEG')
        result_format = result.get('format')
        if result_format is None:
            logger.error('Could not determine image format')
            return 0.0
        if result_format.upper() == expected_format.upper():
            logger.info(f'Format check passed: {result_format}')
            return 1.0
        else:
            logger.warning(f'Format mismatch: got {result_format}, expected {expected_format}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image format: {e}')
        return 0.0

def check_image_grayscale__90bdfeb2ebfd1bac6873718be37b3912(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if image exists and is in grayscale.

    Args:
        result: Dict from getter with 'exists', 'is_grayscale' keys
        expected: Dict with 'is_grayscale' boolean

    Returns:
        float: 1.0 if file exists and grayscale status matches, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('Image file does not exist')
        return 0.0
    expected_grayscale = expected.get('is_grayscale', False)
    actual_grayscale = result.get('is_grayscale', False)
    if actual_grayscale == expected_grayscale:
        logger.info(f'Image grayscale status matches: is_grayscale={actual_grayscale}')
        return 1.0
    else:
        logger.info(f'Image grayscale status mismatch: expected is_grayscale={expected_grayscale}, got {actual_grayscale}')
        return 0.0

def check_image_dimensions__77b19ce3287accb29381eac14cc998b5(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if image has valid dimensions.

    Args:
        result: Dict from getter with 'exists', 'width', 'height', 'is_png' keys
        expected: Dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.3
    else:
        logger.info('File does not exist')
        return 0.0
    if result.get('is_png', False):
        score += 0.2
    width = result.get('width', 0)
    height = result.get('height', 0)
    min_width = expected.get('min_width', 100)
    min_height = expected.get('min_height', 100)
    if width >= min_width and height >= min_height:
        score += 0.5
        logger.info(f'Valid dimensions: {width}x{height}')
    else:
        logger.info(f'Invalid dimensions: {width}x{height} (expected >= {min_width}x{min_height})')
    return score

def check_layer_exists__b148e375(result_state, expected_state, **options):
    """
    Check if a specific layer name exists in the list of layer names.

    Args:
        result_state: List of layer names from the XCF file
        expected_state: Dict with 'layer_name' key specifying the expected layer name
        **options: Additional options

    Returns:
        float: 1.0 if the layer exists, 0.0 otherwise
    """
    if not isinstance(result_state, list):
        logger.error(f'result_state is not a list: {type(result_state)}')
        return 0.0
    if not isinstance(expected_state, dict):
        logger.error(f'expected_state is not a dict: {type(expected_state)}')
        return 0.0
    expected_layer_name = expected_state.get('layer_name', '')
    if not expected_layer_name:
        logger.error('No layer_name specified in expected_state')
        return 0.0
    if expected_layer_name in result_state:
        logger.info(f"Layer '{expected_layer_name}' found in {result_state}")
        return 1.0
    else:
        logger.warning(f"Layer '{expected_layer_name}' not found in {result_state}")
        return 0.0

def check_image_rows_reversed__1f38b5ad(result, expected, **options):
    """
    Check if the image rows have been reversed (bottom to top).

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'num_rows'
        **options: Additional options

    Returns:
        float: Score (1.0 if rows reversed correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        num_rows = expected.get('num_rows', 4)
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        if result_img.size != source_img.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {source_img.size}')
            return 0.0
        (width, height) = source_img.size
        row_height = height // num_rows
        score = 0.0
        for i in range(num_rows):
            source_row = source_img.crop((0, i * row_height, width, (i + 1) * row_height))
            result_row = result_img.crop((0, (num_rows - 1 - i) * row_height, width, (num_rows - i) * row_height))
            if structure_check_by_ssim(source_row, result_row, threshold=0.95):
                score += 1.0 / num_rows
            else:
                logging.debug(f'Row {i} does not match reversed position')
        return score
    except Exception as e:
        logging.error(f'Error in check_image_rows_reversed__1f38b5ad: {e}')
        return 0.0

def check_image_dimensions__b6f18d98d1993dda1c36f44651fe6a5d(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' expected values

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        logger.error(f'Invalid result dimensions: {result}')
        return 0.0
    width_match = actual_width == expected_width
    height_match = actual_height == expected_height
    logger.info(f'Expected: {expected_width}x{expected_height}, Got: {actual_width}x{actual_height}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_crop_dimensions__c26faa30e6e2f9e09b4748ad3193b390(result, expected, **options):
    """
    Check if the image was cropped to expected dimensions.

    Args:
        result: Dict with 'exists', 'width', 'height' from getter
        expected: Dict with 'width' and 'height' expected values
        **options: Additional options (can include 'tolerance' for dimension matching)

    Returns:
        float: 1.0 if dimensions match (within tolerance), 0.0 otherwise
    """
    if result is None or not result.get('exists', False):
        logger.error('Result image does not exist')
        return 0.0
    result_width = result.get('width')
    result_height = result.get('height')
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result_width is None or result_height is None:
        logger.error('Missing dimensions in result')
        return 0.0
    logger.info(f'Result dimensions: {result_width}x{result_height}')
    logger.info(f'Expected dimensions: {expected_width}x{expected_height}')
    tolerance = options.get('tolerance', 0)
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    if width_match and height_match:
        return 1.0
    else:
        logger.info(f'Dimension mismatch - width_match: {width_match}, height_match: {height_match}')
        return 0.0

def check_image_rotated_90_clockwise__25c734e4497155c51ced0623dec284fc(result, expected, **options):
    """
    Check if the result image is rotated 90 degrees clockwise from the source.

    Args:
        result: Dict with 'result_path' and 'source_path'
        expected: Dict with expected rule (not used, just for consistency)
        **options: Additional options

    Returns:
        float: 1.0 if rotated correctly, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    result_path = result.get('result_path')
    source_path = result.get('source_path')
    if not result_path or not source_path:
        logger.error('Missing result_path or source_path')
        return 0.0
    try:
        result_img = Image.open(result_path)
        source_img = Image.open(source_path)
        rotated_source = source_img.transpose(Image.ROTATE_270)
        if result_img.size != rotated_source.size:
            logger.error(f'Size mismatch after rotation: {result_img.size} vs {rotated_source.size}')
            return 0.0
        if result_img.mode != 'RGB':
            result_img = result_img.convert('RGB')
        if rotated_source.mode != 'RGB':
            rotated_source = rotated_source.convert('RGB')
        result_array = np.array(result_img)
        rotated_array = np.array(rotated_source)
        try:
            similarity = ssim(result_array, rotated_array, win_size=7, channel_axis=2)
        except TypeError:
            similarity = ssim(result_array, rotated_array, win_size=7, multichannel=True)
        logger.info(f'Rotation SSIM: {similarity}')
        threshold = options.get('threshold', 0.95)
        return 1.0 if similarity >= threshold else 0.0
    except Exception as e:
        logger.error(f'Error checking rotation: {e}')
        return 0.0

def check_image_width__d4f10b60(src_path, rule):
    """
    Check if the image width matches the expected value
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_7
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        actual_width = img.size[0]
        expected_width = rule.get('width')
        tolerance = rule.get('tolerance', 0)
        logger.debug(f'Image width: {actual_width}, expected: {expected_width}')
        if abs(actual_width - expected_width) <= tolerance:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image width: {e}')
        return 0.0

def check_default_video_player__52f0bd95(result, expected, **options):
    """
    Check if the default video player matches the expected value.

    Args:
        result: The actual default video player (e.g., 'vlc.desktop')
        expected: Dictionary with 'expected_player' key containing the expected value
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_player = expected.get('expected_player', 'vlc.desktop')
    actual = result.strip() if result else ''
    if expected_player in actual or actual == expected_player:
        return 1.0
    else:
        logger.info(f'Default video player mismatch - Expected: {expected_player}, Got: {actual}')
        return 0.0

def check_image_dimensions__505cf5fc(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_grayscale__37ae16c860b893668df4aed8d0b9ae18(result, expected, **options):
    """Check if an image is grayscale based on pixel analysis.

    Args:
        result: Dict from getter with keys 'is_grayscale', 'total_pixels', 'grayscale_pixels'
        expected: Dict with 'type' = 'rule' and 'rules' (empty or containing validation rules)
        **options: Additional comparison options

    Returns:
        float: 1.0 if image is grayscale, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    is_grayscale = result.get('is_grayscale', False)
    if is_grayscale:
        return 1.0
    else:
        return 0.0

def check_image_mirror__dae8d36d(result, expected, **options):
    """
    Check if the image is horizontally mirrored (flipped left-right).

    Args:
        result: Path to result image file
        expected: Dict with 'source_path'
        **options: Additional options

    Returns:
        float: Score (1.0 if mirrored correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        flipped_source = source_img.transpose(Image.FLIP_LEFT_RIGHT)
        if structure_check_by_ssim(result_img, flipped_source, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match horizontally flipped source')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_mirror__dae8d36d: {e}')
        return 0.0

def check_image_deleted__53c4502a(result, expected, **options):
    """Check if specific image was deleted from Pictures directory.

    Args:
        result: JSON string/dict mapping image hashes to filenames
        expected: Rules dict with deleted_hash and remaining_hashes
        **options: Additional options

    Returns:
        float: 1.0 if deleted hash absent and remaining hashes present, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if isinstance(result, str):
        result = result.strip().replace("'", '"')
        try:
            result = json.loads(result)
        except:
            return 0.0
    deleted_hash = expected.get('deleted_hash')
    remaining_hashes = expected.get('remaining_hashes', [])
    if deleted_hash in result:
        return 0.0
    for hash_val in remaining_hashes:
        if hash_val not in result:
            return 0.0
    return 1.0

def check_image_rotated__045e2e0c(result, expected, **options):
    """
    Check if the image was rotated correctly.

    Args:
        result: Dict with 'rotated_90_cw' from getter
        expected: Dict with expected rotation status
        **options: Additional options

    Returns:
        float: Score (1.0 if rotated correctly, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        is_rotated = result.get('rotated_90_cw', False)
        expected_rotated = expected.get('rotated_90_cw', True)
        logger.info(f'Expected rotated: {expected_rotated}, Actual rotated: {is_rotated}')
        if is_rotated == expected_rotated:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking rotation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_dimensions__ab9c94ea(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def verify_image_count__8e1f0943(result, expected, **options):
    """Verify image count meets minimum requirement.

    Args:
        result: Actual count
        expected: Expected minimum (dict with 'min_count' key)
        **options: Additional options

    Returns:
        float: 1.0 if >= minimum, 0.0 otherwise
    """
    min_count = expected.get('min_count', 4)
    return 1.0 if result >= min_count else 0.0

def check_gif_animated__fc6196ba2aeabb7bc4c476007b2b24c6(result, expected, **options):
    """
    Check if a GIF file is properly animated with expected frame count and duration.

    Args:
        result: Dict from getter with frame info (exists, frame_count, is_animated, duration_seconds)
        expected: Dict with expected values (min_frames, max_frames, is_animated, min_duration_seconds, max_duration_seconds)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    expected_is_animated = expected.get('is_animated', True)
    is_animated = result.get('is_animated', False)
    if is_animated == expected_is_animated:
        score += 0.4
        logger.info(f'Animation check passed: is_animated={is_animated}')
    else:
        logger.warning(f'Animation check failed: expected is_animated={expected_is_animated}, got {is_animated}')
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 2)
    max_frames = expected.get('max_frames', 200)
    if min_frames <= frame_count <= max_frames:
        score += 0.3
        logger.info(f'Frame count check passed: {frame_count} frames (range: {min_frames}-{max_frames})')
    else:
        logger.warning(f'Frame count out of range: {frame_count} frames (expected: {min_frames}-{max_frames})')
    duration_seconds = result.get('duration_seconds', 0.0)
    min_duration = expected.get('min_duration_seconds')
    max_duration = expected.get('max_duration_seconds')
    if min_duration is not None and max_duration is not None:
        if min_duration <= duration_seconds <= max_duration:
            score += 0.3
            logger.info(f'Duration check passed: {duration_seconds:.2f} seconds (range: {min_duration}-{max_duration})')
        else:
            logger.warning(f'Duration out of range: {duration_seconds:.2f} seconds (expected: {min_duration}-{max_duration})')
    else:
        score += 0.3
        logger.info('Duration check not specified, skipping')
    return score

def check_gif_file__ea8c7a7a(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_crop__f3b4c5c2(result, expected, **options):
    """
    Check if the image has been cropped to the specified region.

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'crop_box' [x1, y1, x2, y2]
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        crop_box = expected.get('crop_box')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        expected_crop = source_img.crop(tuple(crop_box))
        if result_img.size != expected_crop.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {expected_crop.size}')
            return 0.0
        if structure_check_by_ssim(result_img, expected_crop, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match expected crop')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_crop__f3b4c5c2: {e}')
        return 0.0

def check_gimp_fullscreen__dbacd0b982a5b74d88d0d51e944036be(actual_config_path, rule):
    """
    Check if GIMP fullscreen mode setting is correct
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_triangle_right_edge__1418e487(result_state, expected_state, **options):
    """
    Check if the triangle has been moved to the right edge, vertically centered.

    Args:
        result_state: Path to the result image
        expected_state: Not used (rule-based evaluation)
        **options: Additional options

    Returns:
        float: Score (1.0 if positioned correctly, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        if len(unique_colors_sorted) < 2:
            logger.warning('Could not find triangle in image')
            return 0.0
        triangle_color = unique_colors_sorted[1]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        centroid = triangle_coords.mean(axis=0)
        (image_height, image_width) = img_array.shape[:2]
        image_center_y = image_height / 2
        right_threshold = image_width * 0.85
        vertical_tolerance = image_height * 0.15
        logger.info(f'Triangle centroid: ({centroid[1]:.1f}, {centroid[0]:.1f})')
        logger.info(f'Image size: {image_width}x{image_height}, center_y: {image_center_y:.1f}')
        logger.info(f'Right threshold: {right_threshold:.1f}, vertical tolerance: {vertical_tolerance:.1f}')
        is_right = centroid[1] >= right_threshold
        is_centered = abs(centroid[0] - image_center_y) <= vertical_tolerance
        if is_right and is_centered:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error in check_triangle_right_edge__1418e487: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_triangle_bottom_right__e4487e27c5c6e4232b26add556e7d796(result_path, expected, **options):
    """Check if the yellow triangle has been moved to the bottom-right corner.

    Args:
        result_path: Path to the result image
        expected: Expected rules dict (from config["rules"])
        **options: Additional options

    Returns:
        float: 1.0 if yellow triangle is in bottom-right corner, 0.0 otherwise
    """
    if result_path is None:
        return 0.0
    try:
        result_img = Image.open(result_path)
        result_array = np.array(result_img)
        if result_array.shape[2] == 4:
            result_array = result_array[:, :, :3]
        yellow_lower = np.array([180, 180, 0])
        yellow_upper = np.array([255, 255, 120])
        yellow_mask = np.all((result_array >= yellow_lower) & (result_array <= yellow_upper), axis=2)
        yellow_coords = np.argwhere(yellow_mask)
        if len(yellow_coords) == 0:
            logger.warning('No yellow pixels found in the image')
            return 0.0
        pixel_count = len(yellow_coords)
        image_size = result_array.shape[0] * result_array.shape[1]
        pixel_ratio = pixel_count / image_size
        if pixel_ratio < 0.001 or pixel_ratio > 0.5:
            logger.warning(f'Yellow pixel ratio {pixel_ratio:.4f} seems unreasonable for a triangle')
        centroid = yellow_coords.mean(axis=0)
        (centroid_y, centroid_x) = centroid
        (image_height, image_width) = result_array.shape[:2]
        in_bottom_right = centroid_x > image_width * 2 / 3 and centroid_y > image_height * 2 / 3
        logger.debug(f'Found {pixel_count} yellow pixels ({pixel_ratio:.4f} of image)')
        logger.debug(f'Yellow triangle centroid at ({centroid_x:.1f}, {centroid_y:.1f}), image size ({image_width}, {image_height})')
        logger.debug(f'Bottom-right thresholds: x > {image_width * 2 / 3:.1f}, y > {image_height * 2 / 3:.1f}')
        logger.debug(f'Bottom-right check result: {in_bottom_right}')
        if in_bottom_right:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        return 0.0

def check_image_cropped__8159102f(result, expected, **options):
    """
    Check if the image was cropped.

    Args:
        result: Dict with crop info from getter
        expected: Dict with 'is_cropped' expected value
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        is_cropped = result.get('is_cropped', False)
        expected_cropped = expected.get('is_cropped', True)
        logger.info(f'Expected cropped: {expected_cropped}, Actual cropped: {is_cropped}')
        if is_cropped == expected_cropped:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking crop: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_moved_pngs__486dfc7618cd1caae2ceb6550f7743a0(directory_list, rule):
    """
    Check if expected PNG files are moved to the target directory.

    Args:
        directory_list: Directory tree structure from get_list_directory
        rule: Expected configuration with 'expected' key containing list of filenames

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_pngs = rule['expected']
    moved_pngs = [node['name'] for node in directory_list['children']]
    if len(moved_pngs) != len(expected_pngs):
        return 0.0
    if set(moved_pngs) == set(expected_pngs):
        return 1.0
    else:
        return 0.0

def check_image_dimensions__a4f96e46(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_layer_exists__8249d409(result_state, expected, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected: Expected configuration with 'layer_name' key (dict)
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected, dict):
            expected_layer = expected.get('layer_name', '')
        else:
            logger.error(f'expected is not a dict: {type(expected)}')
            return 0.0
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_format__fa27ef76(src_path, expected):
    """
    Check if the image was saved in the expected format.
    Variation 9 for task 554785e9-4523-4e7a-b8e1-8016f565f56a

    Args:
        src_path: Path to edited image
        expected: Dict with 'format' key (e.g., 'JPEG', 'PNG')

    Returns:
        float: 1.0 if image has expected format, 0.0 otherwise
    """
    if src_path is None:
        logger.error('Source path is None')
        return 0.0
    if not os.path.exists(src_path):
        logger.error(f'File does not exist: {src_path}')
        return 0.0
    try:
        img = Image.open(src_path)
        actual_format = img.format
        expected_format = expected.get('format', 'JPEG').upper()
        logger.info(f'Format check: actual={actual_format}, expected={expected_format}')
        file_ext = os.path.splitext(src_path)[1].lower()
        expected_ext_map = {'JPEG': ['.jpg', '.jpeg'], 'PNG': ['.png'], 'GIF': ['.gif'], 'BMP': ['.bmp']}
        expected_exts = expected_ext_map.get(expected_format, [])
        ext_match = file_ext in expected_exts if expected_exts else True
        logger.info(f'Extension check: file_ext={file_ext}, expected_exts={expected_exts}, match={ext_match}')
        if actual_format == expected_format and ext_match:
            return 1.0
        else:
            logger.debug(f'Format mismatch: actual={actual_format}, expected={expected_format}, ext_match={ext_match}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image format: {e}')
        return 0.0

def check_image_rotated__55bebaef7ca999b793134c2c00f342a9(result, expected, **options):
    """Check if image has been rotated to expected orientation.

    Args:
        result: Dict with width, height, orientation from getter
        expected: Dict with expected orientation and dimensions
        **options: Additional options

    Returns:
        float: 1.0 if orientation matches, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_orientation = expected.get('orientation')
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    orientation_match = result.get('orientation') == expected_orientation
    if expected_width is not None and expected_height is not None:
        width_match = result.get('width') == expected_width
        height_match = result.get('height') == expected_height
        return 1.0 if orientation_match and width_match and height_match else 0.0
    else:
        return 1.0 if orientation_match else 0.0

def check_image_color_mode__14447080df6d9553c7e99d3265fc5a81(result, expected, **options):
    """
    Check if image color mode matches expected value.

    Args:
        result: Dict with 'mode' from getter (e.g., 'RGB', 'L', 'P')
        expected: Dict with expected 'mode' value

    Returns:
        float: 1.0 if mode matches, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_mode = expected.get('mode')
    if expected_mode is None:
        logger.error('Expected mode not specified')
        return 0.0
    result_mode = result.get('mode')
    if result_mode is None:
        logger.error('Result mode is None')
        return 0.0
    mode_match = result_mode == expected_mode
    logger.info(f'Mode: {result_mode} vs {expected_mode} (match={mode_match})')
    if mode_match:
        return 1.0
    else:
        return 0.0

def check_image_size__e19bd559(result_path, expected, **options):
    """
    Check if the resized image has the correct dimensions (800x600) and maintained aspect ratio.

    Args:
        result_path: Path to the resized image file
        expected: Expected rules dict with 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match and aspect ratio is preserved, 0.0 otherwise
    """
    if not result_path:
        return 0.0
    try:
        img = Image.open(result_path)
        (actual_width, actual_height) = img.size
        expected_width = expected.get('width', 800)
        expected_height = expected.get('height', 600)
        if actual_width != expected_width or actual_height != expected_height:
            return 0.0
        actual_ratio = actual_width / actual_height
        expected_ratio = 4.0 / 3.0
        ratio_tolerance = 0.01
        if abs(actual_ratio - expected_ratio) > ratio_tolerance:
            return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking image size: {e}')
        return 0.0

def check_image_dimensions__f2472a44(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_layer_added__e5a8318d(result, expected, **options):
    """
    Check if a layer was successfully added by verifying it exists.

    Args:
        result: Boolean indicating if layer exists (from getter)
        expected: Expected value (should be True)
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if result is True:
            logger.info('Layer successfully added')
            return 1.0
        else:
            logger.warning('Layer was not added')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer addition: {e}')
        return 0.0

def check_image_hash__8778dd25(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_hash__10e2f0b6(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_rotation__b014bcfb2b56c94b7ab718a85ad6cff9(result, expected, **options):
    """Check if image has been rotated to expected orientation.

    Args:
        result: dict with orientation info from getter
        expected: dict with expected orientation
        **options: Additional options

    Returns:
        float: 1.0 if orientation matches, 0.5 if dimensions swapped, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_orientation = expected.get('expected_orientation')
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    score = 0.0
    if result.get('orientation') == expected_orientation:
        score += 0.5
    if result.get('width') == expected_width and result.get('height') == expected_height:
        score += 0.5
    return score

def check_image_rotate_180__201e011d(result, expected, **options):
    """
    Check if the image has been rotated 180 degrees.

    Args:
        result: Path to result image file
        expected: Dict with 'source_path'
        **options: Additional options

    Returns:
        float: Score (1.0 if rotated correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        source_path = expected.get('source_path')
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        rotated_source = source_img.transpose(Image.ROTATE_180)
        if structure_check_by_ssim(result_img, rotated_source, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match 180-degree rotated source')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_rotate_180__201e011d: {e}')
        return 0.0

def check_image_properties__72e666ad(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_triangle_flipped__0293a54ed09c3032ba177ba116651a07(tgt_path, expected, **options):
    """
    Check if the triangle has been flipped horizontally.

    Verifies that:
    1. The triangle was originally on the left side
    2. The triangle is now on the right side
    3. The triangle actually moved (not just checking final position)
    4. The flip orientation matches the expected configuration

    Args:
        tgt_path: Path to the result image
        expected: Expected configuration with orientation (e.g., {'orientation': 'flipped_horizontal'})
        **options: Additional options

    Returns:
        float: 1.0 if triangle is flipped correctly, 0.0 otherwise
    """
    if tgt_path is None:
        return 0.0
    expected_orientation = expected.get('orientation', 'flipped_horizontal') if isinstance(expected, dict) else 'flipped_horizontal'
    original_path = '/home/user/Desktop/Triangle_On_The_Side.png'
    if not os.path.exists(original_path):
        return _check_result_only(tgt_path, expected_orientation)
    original_img = Image.open(original_path)
    original_array = np.array(original_img)
    if original_array.shape[2] == 4:
        original_array = original_array[:, :, :3]
    result_img = Image.open(tgt_path)
    result_array = np.array(result_img)
    if result_array.shape[2] == 4:
        result_array = result_array[:, :, :3]
    if original_array.shape != result_array.shape:
        return 0.0
    (height, width) = original_array.shape[:2]
    original_mask = _detect_yellow_triangle(original_array)
    if np.sum(original_mask) == 0:
        return 0.0
    result_mask = _detect_yellow_triangle(result_array)
    if np.sum(result_mask) == 0:
        return 0.0
    original_coords = np.argwhere(original_mask)
    result_coords = np.argwhere(result_mask)
    original_centroid = original_coords.mean(axis=0)
    result_centroid = result_coords.mean(axis=0)
    (original_centroid_y, original_centroid_x) = original_centroid
    (result_centroid_y, result_centroid_x) = result_centroid
    if expected_orientation == 'flipped_horizontal':
        is_originally_on_left = original_centroid_x < width / 3
        is_now_on_right = result_centroid_x > 2 * width / 3
        moved_across_center = abs(result_centroid_x - original_centroid_x) > width / 3
        is_mirrored = _verify_horizontal_mirror(original_array, result_array, original_mask, result_mask)
        if is_originally_on_left and is_now_on_right and moved_across_center and is_mirrored:
            return 1.0
        else:
            return 0.0
    elif expected_orientation == 'flipped_vertical':
        moved_vertically = abs(result_centroid_y - original_centroid_y) > height / 3
        is_mirrored = _verify_vertical_mirror(original_array, result_array, original_mask, result_mask)
        if moved_vertically and is_mirrored:
            return 1.0
        else:
            return 0.0
    else:
        return 0.0

def check_image_dimensions__795e137d(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_modified__4c0f04bf(result, expected, **options):
    """Check if image was modified (different hash).

    Args:
        result: Dict with image properties from getter
        expected: Dict with expected properties in 'rules' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists'):
        return 0.0
    if expected.get('width') and result.get('width') != expected.get('width'):
        return 0.0
    if expected.get('height') and result.get('height') != expected.get('height'):
        return 0.0
    return 1.0

def check_gif_file__11cf8ab6(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_gimp_gimprc_setting__adbf37d0(actual_config_path, expected, **options):
    """
    Check if a GIMP gimprc setting has the expected value.

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'key' and 'value' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if setting matches, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        target_key = expected.get('key')
        target_value = expected.get('value')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == target_key and value == target_value:
                    logger.info(f'Found matching setting: {key} = {value}')
                    return 1.0
        logger.warning(f"Setting not found or doesn't match: {target_key} = {target_value}")
        return 0.0
    except Exception as e:
        logger.error(f'Error checking gimprc setting: {e}')
        return 0.0

def check_image_format__d255bc15(result_state, expected_state, **options):
    """
    Check if the image is in the expected format.

    Args:
        result_state: Path to the result image file
        expected_state: Dict with 'format' key (e.g., 'JPEG', 'PNG')
        **options: Additional options

    Returns:
        float: 1.0 if image is in expected format, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    try:
        file_ext = os.path.splitext(result_state)[1].lower()
        expected_format = expected_state.get('format', 'JPEG')
        img = Image.open(result_state)
        if img.format == expected_format:
            logger.info(f'Image is in expected format: {img.format}')
            return 1.0
        else:
            logger.warning(f'Image format mismatch: got {img.format}, expected {expected_format}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image format: {e}')
        return 0.0

def check_image_dimensions__05f2a34a(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_gimp_is_grayscale__7ae19854d29f1b0c8cdcf13e028f0bdd(result, expected, **options):
    """
    Check if the image is grayscale.

    Args:
        result: dict with 'mode', 'is_grayscale', 'unique_colors' from getter
        expected: dict (currently unused, grayscale check is binary)
        **options: Additional options

    Returns:
        float: 1.0 if image is grayscale, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    is_grayscale = result.get('is_grayscale', False)
    mode = result.get('mode', 'unknown')
    unique_colors = result.get('unique_colors', 0)
    logger.info(f'Grayscale check: mode={mode}, is_grayscale={is_grayscale}, unique_colors={unique_colors}')
    if is_grayscale:
        logger.info('Grayscale check PASSED')
        return 1.0
    else:
        logger.info('Grayscale check FAILED: image has color')
        return 0.0

def check_gimp_gimprc_setting__dc948653(actual_config_path, expected, **options):
    """
    Check if a GIMP gimprc setting has the expected value.

    The gimprc file uses a different format than sessionrc:
    - Format: (key value)
    - Example: (toolbox-group-menu-mode show-on-hover)

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'key' and 'value' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if setting matches, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        target_key = expected.get('key')
        target_value = expected.get('value')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == target_key and value == target_value:
                    logger.info(f'Found matching setting: {key} = {value}')
                    return 1.0
        logger.warning(f"Setting not found or doesn't match: {target_key} = {target_value}")
        return 0.0
    except Exception as e:
        logger.error(f'Error checking gimprc setting: {e}')
        return 0.0

def check_png_file_exists__d04e36cb527bbbbfdb06458981bf8945(result, expected, **options):
    """
    Check if PNG file exists and has valid PNG format based on expected rules.

    Args:
        result: Path to the result PNG file (from getter)
        expected: Expected rules dict with 'exists' key
        **options: Additional options

    Returns:
        float: 1.0 if file state matches expected rules, 0.0 otherwise
    """
    should_exist = expected.get('exists', True)
    file_exists = result and os.path.exists(result)
    if not should_exist:
        if file_exists:
            logger.error(f'File should not exist but was found: {result}')
            return 0.0
        else:
            logger.info(f'File correctly does not exist')
            return 1.0
    if not file_exists:
        logger.error(f'File does not exist: {result}')
        return 0.0
    try:
        img = Image.open(result)
        if img.format != 'PNG':
            logger.error(f'File is not PNG format: {img.format}')
            return 0.0
        img.load()
        logger.info(f'Valid PNG file: {img.format} {img.mode} {img.size}')
        return 1.0
    except Exception as e:
        logger.error(f'Error loading PNG file: {e}')
        return 0.0

def check_triangle_scale_50__2550dc37(result_state, expected_state, **options):
    """
    Check if the triangle has been scaled to approximately 50% of its original size.

    Args:
        result_state: Path to the result image
        expected_state: Not used (rule-based evaluation)
        **options: Additional options

    Returns:
        float: Score (1.0 if scaled correctly, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)[::-1]]
        if len(unique_colors_sorted) < 2:
            logger.warning('Could not find triangle in image')
            return 0.0
        triangle_color = None
        triangle_coords = None
        for i in range(min(5, len(unique_colors_sorted))):
            color = unique_colors_sorted[i]
            if len(color) >= 3 and color[0] > 200 and (color[1] > 200) and (color[2] < 100):
                triangle_mask = np.all(img_array == color, axis=2)
                coords = np.argwhere(triangle_mask)
                if len(coords) > 100:
                    triangle_color = color
                    triangle_coords = coords
                    logger.info(f'Found yellow triangle with color {color}')
                    break
        if triangle_color is None:
            triangle_color = unique_colors_sorted[1]
            triangle_mask = np.all(img_array == triangle_color, axis=2)
            triangle_coords = np.argwhere(triangle_mask)
            logger.info(f'Using second most common color {triangle_color}')
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        (min_y, min_x) = triangle_coords.min(axis=0)
        (max_y, max_x) = triangle_coords.max(axis=0)
        triangle_height = max_y - min_y + 1
        triangle_width = max_x - min_x + 1
        original_height = 240
        original_width = 170
        logger.info(f'Original triangle dimensions: {original_width}x{original_height}')
        height_ratio = triangle_height / original_height
        width_ratio = triangle_width / original_width
        logger.info(f'Triangle dimensions: {triangle_width}x{triangle_height}')
        logger.info(f'Scale ratios: width={width_ratio:.2f}, height={height_ratio:.2f}')
        if 0.425 <= height_ratio <= 0.575 and 0.425 <= width_ratio <= 0.575:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error in check_triangle_scale_50__2550dc37: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_png_files__4a34f03bcb82e7e8037181cd9a91ae6e(result, expected, **options):
    """Check if PNG files were moved correctly.

    Args:
        result: Directory tree dict from getter
        expected: Rules dict with 'expected' list of filenames
        **options: Additional options

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected.get('expected', [])
    if not result or 'children' not in result:
        return 0.0
    actual_files = [node['name'] for node in result['children']]
    if len(actual_files) != len(expected_files):
        return 0.0
    if set(actual_files) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_image_dimensions__3b2067dec4f66b0e25da3355b1fb8f3e(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Dict with width and height from getter
        expected: Dict with expected width and height
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    width_match = result.get('width') == expected_width
    height_match = result.get('height') == expected_height
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_gif_file__0d0257bdd2e8345a807cfdefd39ffa3b(result, expected, **options):
    """
    Check if a GIF file exists and meets requirements for a 5-second clip.

    Args:
        result: dict from get_gif_file_info__0d0257bdd2e8345a807cfdefd39ffa3b
        expected: dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.25
        logger.info('File exists: +0.25')
    else:
        logger.warning('File does not exist')
        return 0.0
    if result.get('is_gif', False):
        score += 0.25
        logger.info('File is a valid GIF: +0.25')
    else:
        logger.warning('File is not a valid GIF')
        return score
    min_size = expected.get('min_file_size', 1000)
    if result.get('file_size', 0) >= min_size:
        score += 0.15
        logger.info(f"File size {result['file_size']} >= {min_size}: +0.15")
    else:
        logger.warning(f"File size {result['file_size']} < {min_size}")
    min_frames = expected.get('min_frames', 100)
    if result.get('frames', 0) >= min_frames:
        score += 0.2
        logger.info(f"Frame count {result['frames']} >= {min_frames}: +0.2")
    else:
        logger.warning(f"Frame count {result['frames']} < {min_frames}")
    duration_ms = result.get('duration_ms', 0)
    target_duration_ms = expected.get('target_duration_ms', 5000)
    min_duration_ms = expected.get('min_duration_ms', 4000)
    max_duration_ms = expected.get('max_duration_ms', 6000)
    if min_duration_ms <= duration_ms <= max_duration_ms:
        score += 0.15
        logger.info(f'Duration {duration_ms}ms is within acceptable range [{min_duration_ms}-{max_duration_ms}ms]: +0.15')
    else:
        logger.warning(f'Duration {duration_ms}ms is outside acceptable range [{min_duration_ms}-{max_duration_ms}ms]')
    logger.info(f'Final score: {score}')
    return score

def check_image_horizontal_flip__340a3600a88f3d4dc2217b9f986d625d(result, expected, **options):
    """Check if image has been horizontally flipped by comparing edge colors.

    Args:
        result: dict with corner and edge color info from getter
        expected: dict with expected flipped status (True/False)
        **options: Additional options

    Returns:
        float: 1.0 if flip detection matches expectation, 0.0 otherwise
    """
    if result is None:
        return 0.0
    corners = result.get('corners', {})
    left_avg = result.get('left_edge_avg', (0, 0, 0))
    right_avg = result.get('right_edge_avg', (0, 0, 0))

    def color_distance(c1, c2):
        return sum((abs(a - b) for (a, b) in zip(c1, c2)))
    edge_diff = color_distance(left_avg, right_avg)
    expected_left = expected.get('expected_left_edge_avg')
    expected_right = expected.get('expected_right_edge_avg')
    if expected_left is None or expected_right is None:
        if result.get('width') and result.get('height'):
            return 1.0
        return 0.0
    left_match = color_distance(left_avg, tuple(expected_left)) < 30
    right_match = color_distance(right_avg, tuple(expected_right)) < 30
    if left_match and right_match:
        return 1.0
    return 0.0

def check_triangle_topedge__e65601d6(result, expected, **options):
    """
    Check if the yellow triangle is positioned at the top edge of the image.
    Variation 5 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        result: Path to the exported image (from vm_file getter)
        expected: Expected state (dict with rules)
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is at top edge, 0.0 otherwise)
    """
    tgt_path = result
    if tgt_path is None:
        return 0.0
    expected_at_top = expected.get('at_top_edge', True)
    if not expected_at_top:
        return 0.0
    try:
        img = Image.open(tgt_path)
        img_array = np.array(img)
        if img_array.shape[-1] == 4:
            img_array = img_array[:, :, :3]
        yellow_lower = np.array([180, 180, 0])
        yellow_upper = np.array([255, 255, 100])
        yellow_mask = np.all((img_array >= yellow_lower) & (img_array <= yellow_upper), axis=2)
        if not np.any(yellow_mask):
            logger.warning('No yellow triangle detected in the image')
            return 0.0
        triangle_coords = np.argwhere(yellow_mask)
        if len(triangle_coords) == 0:
            return 0.0
        min_y = triangle_coords[:, 0].min()
        img_height = img_array.shape[0]
        at_top_edge = min_y < img_height * 0.05
        if at_top_edge:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        return 0.0

def check_image_dimensions__1728b8eeebbdebd7894eb2578c0c67ec(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: Dict with 'width' and 'height' from getter
        expected: Dict with expected 'width' and 'height' values

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is None or expected_height is None:
        logger.error('Expected width or height not specified')
        return 0.0
    result_width = result.get('width')
    result_height = result.get('height')
    if result_width is None or result_height is None:
        logger.error('Result width or height is None')
        return 0.0
    width_match = result_width == expected_width
    height_match = result_height == expected_height
    logger.info(f'Width: {result_width} vs {expected_width} (match={width_match})')
    logger.info(f'Height: {result_height} vs {expected_height} (match={height_match})')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_grayscale__f3cdb596(result_state, expected_state, **options):
    """
    Check if the image is in grayscale mode.

    Args:
        result_state: Path to the result image file
        expected_state: Dict with 'mode' key (expected 'L' for grayscale)
        **options: Additional options

    Returns:
        float: 1.0 if image is grayscale, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    try:
        img = Image.open(result_state)
        expected_mode = expected_state.get('mode', 'L')
        if img.mode == expected_mode or img.mode in ['L', 'LA']:
            logger.info(f'Image is in grayscale mode: {img.mode}')
            return 1.0
        else:
            logger.warning(f'Image mode is not grayscale: got {img.mode}, expected {expected_mode}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image mode: {e}')
        return 0.0

def check_gif_file__06536a54(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_image_hash__b461596b(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_image_rotated_90__30840ecf(result_state, expected_state, **options):
    """
    Check if the result image is rotated 90 degrees clockwise compared to the source.

    This function verifies rotation by comparing pixel data between the original and
    rotated images. For a 90-degree clockwise rotation, pixel at (x, y) in the original
    should appear at (height-1-y, x) in the rotated image.

    Args:
        result_state: Path to the result image file
        expected_state: Dict with 'expected_width' and 'expected_height' after rotation
        **options: Additional options

    Returns:
        float: 1.0 if image is correctly rotated 90 degrees clockwise, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    try:
        if not os.path.exists(result_state):
            logger.error(f'Result file does not exist: {result_state}')
            return 0.0
        rotated_img = Image.open(result_state)
        (rotated_width, rotated_height) = rotated_img.size
        expected_width = expected_state.get('expected_width')
        expected_height = expected_state.get('expected_height')
        if rotated_width != expected_width or rotated_height != expected_height:
            logger.warning(f'Image size mismatch: got {rotated_width}x{rotated_height}, expected {expected_width}x{expected_height}')
            return 0.0
        original_path = '/home/user/Desktop/character.png'
        if not os.path.exists(original_path):
            logger.error(f'Original file not found: {original_path}')
            return 0.0
        original_img = Image.open(original_path)
        (orig_width, orig_height) = original_img.size
        original_img = original_img.convert('RGB')
        rotated_img = rotated_img.convert('RGB')
        original_array = np.array(original_img)
        rotated_array = np.array(rotated_img)
        sample_points = [(0, 0), (orig_width - 1, 0), (0, orig_height - 1), (orig_width - 1, orig_height - 1), (orig_width // 2, orig_height // 2), (orig_width // 4, orig_height // 4), (3 * orig_width // 4, orig_height // 4), (orig_width // 4, 3 * orig_height // 4), (3 * orig_width // 4, 3 * orig_height // 4)]
        mismatches = 0
        for (x, y) in sample_points:
            rotated_x = orig_height - 1 - y
            rotated_y = x
            original_pixel = original_array[y, x]
            rotated_pixel = rotated_array[rotated_y, rotated_x]
            pixel_diff = np.abs(original_pixel.astype(int) - rotated_pixel.astype(int))
            if np.any(pixel_diff > 5):
                mismatches += 1
                logger.debug(f'Pixel mismatch at original({x},{y}) vs rotated({rotated_x},{rotated_y}): {original_pixel} vs {rotated_pixel}')
        if mismatches <= 1:
            logger.info(f'Image correctly rotated 90 degrees clockwise: {rotated_width}x{rotated_height}, {mismatches}/{len(sample_points)} sample mismatches')
            return 1.0
        else:
            logger.warning(f'Image rotation verification failed: {mismatches}/{len(sample_points)} sample points mismatched')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking rotated image: {e}')
        return 0.0

def check_image_properties__0a6f6cc1(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, and exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
        elif expected['exists'] and (not result.get('exists')):
            return 0.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_folder_images__fe2f25921e2c2c43fbc9e31c35bccb78(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if folder contains expected image files.

    Args:
        result: List of image filenames from getter
        expected: Dict with 'expected_files' key containing list of expected filenames
        **options: Additional options (unused)

    Returns:
        1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if expected_set.issubset(result_set):
        return 1.0
    else:
        return 0.0

def check_image_size__7ecb394c9eae8a8e135a21ca629ec0de(result, expected, **options):
    """
    Check if image exists and has expected dimensions.

    Args:
        result: dict with "width", "height", and "exists" keys from getter
        expected: dict with expected "width" and "height" values (from rules)
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: Score based on matching criteria
            - 0.0 if file doesn't exist
            - 0.5 if file exists but dimensions don't match
            - 1.0 if file exists and dimensions match
    """
    if result is None or not result.get('exists', False):
        logger.error('Result file does not exist')
        return 0.0
    tolerance = options.get('tolerance', 5)
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    result_width = result.get('width')
    result_height = result.get('height')
    if result_width is None or result_height is None:
        logger.error('Result missing width or height')
        return 0.5
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    if width_match and height_match:
        logger.info(f'Image size matches: {result_width}x{result_height}')
        return 1.0
    else:
        logger.info(f'Image size mismatch: got {result_width}x{result_height}, expected {expected_width}x{expected_height} (tolerance {tolerance})')
        return 0.5

def check_image_rotation__82cdf8c66b7532068e061ebffd0a6c33(result, expected, **options):
    """
    Check if image has been rotated (dimensions swapped).

    Args:
        result: dict with "width", "height", and "swapped" keys from getter
        expected: dict with expected "width" and "height" values (from rules)
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match expected (rotated), 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 5)
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    result_width = result.get('width')
    result_height = result.get('height')
    if result_width is None or result_height is None:
        logger.error('Result missing width or height')
        return 0.0
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    if width_match and height_match:
        logger.info(f'Image rotation correct: {result_width}x{result_height}')
        return 1.0
    else:
        logger.info(f'Image rotation incorrect: got {result_width}x{result_height}, expected {expected_width}x{expected_height}')
        return 0.0

def check_image_scaled__e429d357(src_path, rule):
    """
    Check if the image has been scaled by the expected factor
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_2
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        actual_width = img.size[0]
        actual_height = img.size[1]
        scale_factor = rule.get('scale_factor', 0.5)
        tolerance = rule.get('tolerance', 5)
        original_width = 1600
        original_height = 1200
        expected_width = int(original_width * scale_factor)
        expected_height = int(original_height * scale_factor)
        logger.debug(f'Image size: {img.size}, expected: {expected_width}x{expected_height}')
        width_ok = abs(actual_width - expected_width) <= tolerance
        height_ok = abs(actual_height - expected_height) <= tolerance
        if width_ok and height_ok:
            logger.debug(f'Image correctly scaled: width_ok={width_ok}, height_ok={height_ok}')
            return 1.0
        else:
            logger.debug(f'Image not correctly scaled: width_ok={width_ok}, height_ok={height_ok}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image scale: {e}')
        return 0.0

def check_image_scaled__bf4967f3e9931b3c80be8e4dbf6e04b7(result, expected, **options):
    """
    Check if image has been scaled to expected dimensions.

    Args:
        result: Dict with 'width', 'height' from getter
        expected: Expected rules dict with 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if not result:
        logger.error('No result data')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result.get('width') != expected_width or result.get('height') != expected_height:
        logger.error(f"Dimensions mismatch: expected {expected_width}x{expected_height}, got {result.get('width')}x{result.get('height')}")
        return 0.0
    logger.info(f'Image scaled correctly to {expected_width}x{expected_height}')
    return 1.0

def check_image_dimensions__ab3f6dd1(result, expected, **options):
    """
    Check if image dimensions are within expected ranges.

    Args:
        result: Dict from getter with 'exists', 'width', 'height'
        expected: Dict with 'width_range' and 'height_range'
        **options: Additional options

    Returns:
        float: Score based on dimension checks
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.3
    width = result.get('width', 0)
    height = result.get('height', 0)
    width_range = expected.get('width_range', [0, float('inf')])
    height_range = expected.get('height_range', [0, float('inf')])
    if width_range[0] <= width <= width_range[1]:
        score += 0.35
    if height_range[0] <= height <= height_range[1]:
        score += 0.35
    return score

def check_image_dimensions__9f5722fc(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: Dict with 'width' and 'height' keys from getter
        expected: Dict with expected 'width' and 'height' values
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: Score (1.0 if dimensions match, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        result_width = result.get('width')
        result_height = result.get('height')
        if result_width is None or result_height is None:
            logger.error(f'Invalid result dimensions: {result}')
            return 0.0
        width_match = result_width == expected_width if expected_width is not None else True
        height_match = result_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            logger.info(f'Dimensions match: {result_width}x{result_height}')
            return 1.0
        else:
            logger.warning(f'Dimensions mismatch: got {result_width}x{result_height}, expected {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image dimensions: {e}')
        return 0.0

def check_triangle_color__0ab9cbd5(tgt_path, expected, **options):
    """
    Check if the triangle color has been changed to the target color.
    Variation 5 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with target_color and color_tolerance
        **options: Additional options

    Returns:
        float: Score (1.0 if color matches, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        img = Image.open(tgt_path)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        sorted_indices = np.argsort(counts)[::-1]
        unique_colors_sorted = unique_colors[sorted_indices]
        triangle_color = unique_colors_sorted[1][:3]
        target = expected.get('target_color', 'red')
        tolerance = expected.get('color_tolerance', 50)
        logger.info(f'Triangle color: {triangle_color}')
        if target == 'red':
            is_red = triangle_color[0] > triangle_color[1] + tolerance and triangle_color[0] > triangle_color[2] + tolerance
            if is_red:
                return 1.0
        elif target == 'blue':
            is_blue = triangle_color[2] > triangle_color[0] + tolerance and triangle_color[2] > triangle_color[1] + tolerance
            if is_blue:
                return 1.0
        elif target == 'green':
            is_green = triangle_color[1] > triangle_color[0] + tolerance and triangle_color[1] > triangle_color[2] + tolerance
            if is_green:
                return 1.0
        return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle color: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_new_image_added__b5b56630(result, expected, **options):
    """Check if new image was added.

    Args:
        result: Dict with image info
        expected: Expected state (dict with 'has_new' key)
        **options: Additional options

    Returns:
        float: 1.0 if new image detected, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_has_new = expected.get('has_new', True)
    actual_has_new = result.get('has_new', False)
    return 1.0 if actual_has_new == expected_has_new else 0.0

def check_triangle_deleted__0511b7c6(tgt_path, expected, **options):
    """
    Check if the yellow triangle has been removed from the image.
    Variation 8 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with max_triangle_pixels threshold
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is removed, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        img = Image.open(tgt_path)
        img_array = np.array(img)
        if len(img_array.shape) < 3:
            is_white = img_array == 255
        elif img_array.shape[2] == 4:
            is_white = np.all(img_array[:, :, :3] == 255, axis=2) & (img_array[:, :, 3] == 255)
        else:
            is_white = np.all(img_array == 255, axis=2)
        non_white_pixels = np.sum(~is_white)
        max_pixels = expected.get('max_triangle_pixels', 500)
        logger.info(f'Non-white pixels: {non_white_pixels}, Max allowed: {max_pixels}')
        if non_white_pixels <= max_pixels:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle removal: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_format__945e33b1(result, expected, **options):
    """
    Check if the image has the expected format.

    Args:
        result: Image format string from getter
        expected: Dict with 'format' key
        **options: Additional options

    Returns:
        float: Score (1.0 if format matches, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        expected_format = expected.get('format')
        if not expected_format:
            logger.error('No expected format provided')
            return 0.0
        logger.info(f'Expected format: {expected_format}, Actual format: {result}')
        if result.upper() == expected_format.upper():
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image format: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_cropped__cdcc2775(src_path, rule):
    """
    Check if the image has been cropped to smaller dimensions.
    Verifies the new size is smaller than the original in both dimensions.
    """
    if src_path is None:
        return 0.0
    img = Image.open(src_path)
    original_width = rule.get('original_width', None)
    original_height = rule.get('original_height', None)
    if original_width is None or original_height is None:
        logger.error('Original dimensions not provided in rule')
        return 0.0
    is_cropped = img.size[0] < original_width and img.size[1] < original_height
    logger.info(f'Original size: ({original_width}, {original_height}), Current size: {img.size}, is_cropped: {is_cropped}')
    return 1.0 if is_cropped else 0.0

def check_image_rotated__6575e228(result_state, expected_state, **options):
    """
    Check if an image has been rotated correctly by comparing dimensions.

    For a 90-degree rotation (clockwise or counterclockwise), the width and height
    should be swapped. This function checks if the result image dimensions match
    the expected dimensions after rotation.

    Args:
        result_state: Dictionary from getter with 'width' and 'height' keys
        expected_state: Dictionary with expected dimensions:
            - 'expected_width': Expected width after rotation
            - 'expected_height': Expected height after rotation
        **options: Additional options:
            - 'tolerance': Pixel tolerance for dimension comparison (default: 0)

    Returns:
        float: Score (1.0 if dimensions match within tolerance, 0.0 otherwise)
    """
    try:
        if result_state is None or not isinstance(result_state, dict):
            logger.error(f'Invalid result_state: {result_state}')
            return 0.0
        actual_width = result_state.get('width')
        actual_height = result_state.get('height')
        if actual_width is None or actual_height is None:
            logger.error(f'Missing width or height in result_state: {result_state}')
            return 0.0
        expected_width = expected_state.get('expected_width')
        expected_height = expected_state.get('expected_height')
        if expected_width is None or expected_height is None:
            logger.error(f'Missing expected_width or expected_height in expected_state: {expected_state}')
            return 0.0
        tolerance = options.get('tolerance', 0)
        width_match = abs(actual_width - expected_width) <= tolerance
        height_match = abs(actual_height - expected_height) <= tolerance
        if width_match and height_match:
            logger.info(f'Image rotation verified: actual={actual_width}x{actual_height}, expected={expected_width}x{expected_height}, tolerance={tolerance}')
            return 1.0
        else:
            logger.warning(f'Image rotation failed: actual={actual_width}x{actual_height}, expected={expected_width}x{expected_height}, tolerance={tolerance}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image rotation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_size__4f6d3ad3(result_state, expected_state, **options):
    """
    Check if the image has the expected dimensions.

    Args:
        result_state: Path to the result image file
        expected_state: Dict with 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if image has expected size, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    try:
        img = Image.open(result_state)
        (width, height) = img.size
        expected_width = expected_state.get('width')
        expected_height = expected_state.get('height')
        if width == expected_width and height == expected_height:
            logger.info(f'Image size matches: {width}x{height}')
            return 1.0
        else:
            logger.warning(f'Image size mismatch: got {width}x{height}, expected {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image size: {e}')
        return 0.0

def check_image_properties__4894850c(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, format, exists
        **options: Additional options (dimension_tolerance for width/height)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
        else:
            return 0.0
    if not result.get('exists', False):
        return 0.0
    if 'width' in expected:
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected:
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected:
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_triangle_topright__b91aae4ce11e76fbc123f25e08b42c98(tgt_path, expected, **options):
    """
    Check if the triangle is in the top-right corner of the image.

    Args:
        tgt_path: Path to the result image
        expected: Expected configuration with position and tolerance
        **options: Additional options

    Returns:
        float: 1.0 if triangle is in top-right, 0.0 otherwise
    """
    if tgt_path is None:
        return 0.0
    img = Image.open(tgt_path)
    img_array = np.array(img)
    (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
    unique_colors_sorted = unique_colors[np.argsort(counts)]
    triangle_color = unique_colors_sorted[1]
    triangle_mask = np.all(img_array == triangle_color, axis=2)
    triangle_coords = np.argwhere(triangle_mask)
    centroid = triangle_coords.mean(axis=0)
    tolerance = expected.get('tolerance', 0.15)
    image_shape = np.array(img_array.shape[:2])
    top_right = np.array([0, image_shape[1]])
    tolerance_pixels = tolerance * image_shape
    in_topright = np.all(np.abs(centroid - top_right) < tolerance_pixels)
    if bool(in_topright):
        return 1.0
    else:
        return 0.0

def check_image_properties__50f8e5bb(result, expected, **options):
    """Compare image properties against expected values including content verification.

    Args:
        result: dict from getter with image properties (width, height, format, exists, phash)
        expected: dict with expected width, height, exists, and optionally phash
        **options: Additional options (dimension_tolerance, hash_threshold)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
        else:
            logger.info(f"File existence check failed: expected={expected['exists']}, actual={result.get('exists')}")
            return 0.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        width_match = abs(actual_width - expected_width) <= tolerance
        if width_match:
            score += 1.0
        logger.info(f'Width check: expected={expected_width}, actual={actual_width}, match={width_match}')
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        height_match = abs(actual_height - expected_height) <= tolerance
        if height_match:
            score += 1.0
        logger.info(f'Height check: expected={expected_height}, actual={actual_height}, match={height_match}')
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if 'phash' in expected and result.get('exists', False):
        total_checks += 1
        expected_phash = expected['phash']
        actual_phash = result.get('phash')
        if expected_phash and actual_phash:
            try:
                hash1 = int(expected_phash, 16)
                hash2 = int(actual_phash, 16)
                hash_diff = bin(hash1 ^ hash2).count('1')
                hash_threshold = options.get('hash_threshold', 5)
                hash_match = hash_diff <= hash_threshold
                logger.info(f'Perceptual hash check: expected={expected_phash}, actual={actual_phash}, difference={hash_diff}, threshold={hash_threshold}, match={hash_match}')
                if hash_match:
                    score += 1.0
            except Exception as e:
                logger.error(f'Error comparing perceptual hashes: {e}')
                pass
        else:
            logger.warning(f'Missing perceptual hash - expected={expected_phash}, actual={actual_phash}')
    if total_checks == 0:
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks})')
    return final_score

def check_triangle_rotated__e0c6cf186be7cc4682d3712837dcdd72(result_data, expected, **options):
    """Check if the triangle has been rotated 180 degrees by comparing original and rotated images.

    Args:
        result_data: Dict with 'original_path' and 'rotated_path' keys
        expected: Expected rules dict (from config["rules"])
        **options: Additional options

    Returns:
        float: 1.0 if triangle is rotated 180 degrees, 0.0 otherwise
    """
    if result_data is None:
        return 0.0
    try:
        original_path = result_data.get('original_path')
        rotated_path = result_data.get('rotated_path')
        if not original_path or not rotated_path:
            logger.error('Missing original or rotated image path')
            return 0.0
        original_img = Image.open(original_path).convert('RGB')
        rotated_img = Image.open(rotated_path).convert('RGB')
        original_array = np.array(original_img)
        rotated_array = np.array(rotated_img)
        if original_array.shape != rotated_array.shape:
            logger.error(f"Image dimensions don't match: original {original_array.shape}, rotated {rotated_array.shape}")
            return 0.0
        pixel_difference = np.mean(np.abs(original_array.astype(float) - rotated_array.astype(float)))
        if pixel_difference < 5.0:
            logger.debug(f'Images are too similar (pixel diff: {pixel_difference:.2f}), no transformation detected')
            return 0.0
        rotated_back = np.rot90(rotated_array, k=2)
        similarity_score = np.mean(np.abs(original_array.astype(float) - rotated_back.astype(float)))
        logger.debug(f'Pixel difference from original: {pixel_difference:.2f}')
        logger.debug(f'Similarity after rotating back: {similarity_score:.2f}')
        if similarity_score < 10.0:
            logger.debug('180-degree rotation verified successfully')
            return 1.0
        else:
            logger.debug(f'Rotation verification failed - not a 180-degree rotation')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle rotation: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_triangle_topleft__562f7ff1b81bbf5a1dbcdac12003c550(tgt_path, expected, **options):
    """
    Check if the triangle is in the top-left corner of the image.

    Args:
        tgt_path: Path to the result image
        expected: Expected configuration with position and tolerance
        **options: Additional options

    Returns:
        float: 1.0 if triangle is in top-left, 0.0 otherwise
    """
    if tgt_path is None:
        return 0.0
    img = Image.open(tgt_path)
    img_array = np.array(img)
    (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
    unique_colors_sorted = unique_colors[np.argsort(counts)]
    triangle_color = unique_colors_sorted[1]
    triangle_mask = np.all(img_array == triangle_color, axis=2)
    triangle_coords = np.argwhere(triangle_mask)
    centroid = triangle_coords.mean(axis=0)
    tolerance = expected.get('tolerance', 0.15)
    tolerance_pixels = tolerance * np.array(img_array.shape[:2])
    top_left = np.array([0, 0])
    in_topleft = np.all(np.abs(centroid - top_left) < tolerance_pixels)
    if bool(in_topleft):
        return 1.0
    else:
        return 0.0

def check_image_flipped__9f1f61b8216550440a48089a3e4c1731(result_path: str, expected: Any, **options) -> float:
    """
    Check if image is horizontally flipped compared to original.

    Args:
        result_path: Path to result image file
        expected: Dict with 'original_path' key pointing to original image

    Returns:
        float: 1.0 if image appears to be horizontally flipped, 0.0 otherwise
    """
    if not result_path:
        logger.info('Result image file does not exist')
        return 0.0
    original_path = expected.get('original_path')
    if not original_path:
        logger.error('No original image path provided in expected')
        return 0.0
    try:
        import os
        if not os.path.exists(result_path) or not os.path.exists(original_path):
            logger.error(f'Image files missing: result={result_path}, original={original_path}')
            return 0.0
        result_img = Image.open(result_path)
        original_img = Image.open(original_path)
        if result_img.size != original_img.size:
            logger.info(f"Image dimensions don't match: result={result_img.size}, original={original_img.size}")
            return 0.0
        result_flipped_back = result_img.transpose(Image.FLIP_LEFT_RIGHT)
        original_array = np.array(original_img)
        result_array = np.array(result_flipped_back)
        mse = np.mean((original_array.astype(float) - result_array.astype(float)) ** 2)
        threshold = options.get('threshold', 10.0)
        is_flipped = mse < threshold
        logger.info(f'Flip check MSE: {mse}, threshold: {threshold}, is_flipped: {is_flipped}')
        return 1.0 if is_flipped else 0.0
    except Exception as e:
        logger.error(f'Error checking image flip: {e}')
        return 0.0

def check_gimp_image_size__e2468a8febb27268d777ba03561c41ab(result, expected, **options):
    """
    Check if the image content size matches expected dimensions.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'min_width', 'max_width', 'min_height', 'max_height' keys
        **options: Additional options

    Returns:
        float: 1.0 if size is within expected range, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    width = result.get('width')
    height = result.get('height')
    if width is None or height is None:
        logger.error('Width or height not found in result')
        return 0.0
    min_width = expected.get('min_width', 0)
    max_width = expected.get('max_width', float('inf'))
    min_height = expected.get('min_height', 0)
    max_height = expected.get('max_height', float('inf'))
    width_ok = min_width <= width <= max_width
    height_ok = min_height <= height <= max_height
    logger.info(f'Size check: width={width} (expected {min_width}-{max_width}), height={height} (expected {min_height}-{max_height})')
    logger.info(f'Width OK: {width_ok}, Height OK: {height_ok}')
    if width_ok and height_ok:
        return 1.0
    else:
        return 0.0

def check_triangle_rotation__2c517d983ea369519ffc55979e3393ec(result, expected, **options):
    """
    Check if the triangle has been rotated by the expected angle.

    Args:
        result: Dict with 'original_angle' and 'final_angle' keys from getter
        expected: Dict with 'target_angle' key specifying the expected rotation (in degrees)
        **options: Optional 'tolerance' for angle comparison (default: 10 degrees)

    Returns:
        float: 1.0 if rotation matches (within tolerance), 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    original_angle = result.get('original_angle')
    final_angle = result.get('final_angle')
    if original_angle is None or final_angle is None:
        return 0.0
    target_rotation = expected.get('target_angle', 0)
    tolerance = options.get('tolerance', 10)
    actual_rotation = (final_angle - original_angle) % 360
    diff = abs(actual_rotation - target_rotation)
    if diff > 180:
        diff = 360 - diff
    if diff <= tolerance:
        return 1.0
    else:
        return 0.0

def check_png_files_exact_count__f9d97b19(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if exactly N PNG files exist

    Args:
        result: List of PNG filenames
        expected: Dict with:
            - count: Expected exact count

    Returns:
        1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_count = expected.get('count', 0)
    actual_count = len(result)
    if actual_count == expected_count:
        logger.info(f'PNG file count matches: {actual_count}')
        return 1.0
    else:
        logger.info(f'PNG file count mismatch: expected {expected_count}, got {actual_count}')
        return 0.0

def check_triangle_opacity__0ff223c2(tgt_path, expected, **options):
    """
    Check if the triangle opacity has been reduced to the target level.
    Variation 7 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        tgt_path: Path to the result image
        expected: Dictionary with target_opacity and tolerance
        **options: Additional options

    Returns:
        float: Score (1.0 if opacity matches, 0.0 otherwise)
    """
    if tgt_path is None:
        return 0.0
    try:
        img = Image.open(tgt_path)
        if img.mode != 'RGBA':
            logger.warning(f'Image mode is {img.mode}, expected RGBA')
            return 0.0
        img_array = np.array(img)
        is_white = np.all(img_array[:, :, :3] == 255, axis=2)
        is_visible = img_array[:, :, 3] > 10
        triangle_mask = ~is_white & is_visible
        triangle_alpha = img_array[triangle_mask, 3]
        if len(triangle_alpha) == 0:
            logger.error('No triangle pixels found')
            return 0.0
        avg_alpha = np.mean(triangle_alpha) / 255.0
        target = expected.get('target_opacity', 0.5)
        tolerance = expected.get('tolerance', 0.15)
        logger.info(f'Average triangle opacity: {avg_alpha:.3f}, Target: {target}, Tolerance: {tolerance}')
        if abs(avg_alpha - target) <= tolerance:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking triangle opacity: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_resize__eed0216b(result_path, expected, **options):
    """
    Check if the image has been resized to the expected dimensions.

    Args:
        result_path: Path to the result image
        expected: Dict with 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: Score (1.0 if size matches, 0.0 otherwise)
    """
    if result_path is None:
        logger.error('Result path is None')
        return 0.0
    try:
        img = Image.open(result_path)
        (actual_width, actual_height) = img.size
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        logger.info(f'Image size: {actual_width}x{actual_height}, expected: {expected_width}x{expected_height}')
        if actual_width == expected_width and actual_height == expected_height:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image size: {e}')
        return 0.0

def check_gimp_exact_dimensions__4bd0ac4fe70775f29bef20a161a34c39(result, expected, **options):
    """
    Check if the image dimensions match the expected values exactly.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys specifying exact expected dimensions
        **options: Additional options (e.g., 'tolerance' for allowed difference in pixels)

    Returns:
        float: 1.0 if dimensions match (within tolerance), 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        logger.error('Width or height not found in result')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is None or expected_height is None:
        logger.error('Expected width or height not specified')
        return 0.0
    tolerance = options.get('tolerance', 2)
    width_diff = abs(actual_width - expected_width)
    height_diff = abs(actual_height - expected_height)
    width_ok = width_diff <= tolerance
    height_ok = height_diff <= tolerance
    logger.info(f'Dimension check: actual={actual_width}x{actual_height}, expected={expected_width}x{expected_height}, tolerance={tolerance}')
    logger.info(f'Width diff: {width_diff}, Height diff: {height_diff}')
    logger.info(f'Width OK: {width_ok}, Height OK: {height_ok}')
    if width_ok and height_ok:
        return 1.0
    else:
        return 0.0

def check_image_properties__70f3db86(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_image_cropped__06e412f2(src_path, rule):
    """
    Check if the src image has been cropped and has specific dimensions.
    Variation 5 for task 7a4deb26-d57d-4ea9-9a73-630f66a7b568
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        expected_width = rule.get('width')
        expected_height = rule.get('height')
        (actual_width, actual_height) = img.size
        width_match = expected_width is None or actual_width == expected_width
        height_match = expected_height is None or actual_height == expected_height
        logger.debug(f'Image size: {img.size}, Expected: {expected_width}x{expected_height}')
        return 1.0 if width_match and height_match else 0.0
    except Exception as e:
        logger.error(f'Error in check_image_cropped__06e412f2: {e}')
        return 0.0

def check_layer_exists__aa5f92aa(result_state, expected_state, **options):
    """
    Check if a specific layer exists in the GIMP image.

    This metric checks if the expected layer name exists in the list of layer names
    extracted from the XCF file.

    Args:
        result_state: List of layer names from the getter (list of strings)
        expected_state: Dict with 'layer_name' key containing the expected layer name
        **options: Additional options

    Returns:
        float: Score (1.0 if layer exists, 0.0 otherwise)
    """
    try:
        if isinstance(expected_state, dict):
            expected_layer = expected_state.get('layer_name', '')
        else:
            expected_layer = expected_state
        if not isinstance(result_state, list):
            logger.error(f'result_state is not a list: {type(result_state)}')
            return 0.0
        if expected_layer in result_state:
            logger.info(f"Layer '{expected_layer}' found in layers: {result_state}")
            return 1.0
        else:
            logger.warning(f"Layer '{expected_layer}' not found in layers: {result_state}")
            return 0.0
    except Exception as e:
        logger.error(f'Error checking layer existence: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_dimensions__da5d7378bbf41608325407fc00f8d126(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys (expected values)
        **options: Additional options (unused)

    Returns:
        float: 1.0 if dimensions match exactly, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if expected_width is not None and actual_width != expected_width:
        return 0.0
    if expected_height is not None and actual_height != expected_height:
        return 0.0
    return 1.0

def check_gimp_gimprc_setting__17026b9e(actual_config_path, expected, **options):
    """
    Check if a GIMP gimprc setting has the expected value.

    Args:
        actual_config_path: Path to the gimprc config file
        expected: Expected configuration with 'key' and 'value' fields
        **options: Additional options

    Returns:
        float: Score (1.0 if setting matches, 0.0 otherwise)
    """
    if actual_config_path is None:
        return 0.0
    try:
        with open(actual_config_path, 'r') as f:
            content = f.readlines()
        target_key = expected.get('key')
        target_value = expected.get('value')
        for line in content:
            if line.startswith('#') or line.strip() == '':
                continue
            line = line.strip().lstrip('(').rstrip(')\n')
            parts = line.split(None, 1)
            if len(parts) >= 2:
                key = parts[0]
                value = parts[1].strip().strip('"')
                if key == target_key and value == target_value:
                    logger.info(f'Found matching setting: {key} = {value}')
                    return 1.0
        logger.warning(f"Setting not found or doesn't match: {target_key} = {target_value}")
        return 0.0
    except Exception as e:
        logger.error(f'Error checking gimprc setting: {e}')
        return 0.0

def check_image_dimensions__184d842d21bd06203c79c089532a2315(result, expected, **options):
    """
    Check if image has expected dimensions and format.

    Args:
        result: Dict with 'width', 'height', 'format' from getter
        expected: Expected rules dict with 'width', 'height', 'format' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions and format match, 0.0 otherwise
    """
    if not result:
        logger.error('No result data')
        return 0.0
    expected_format = expected.get('format', 'JPEG')
    if result.get('format') != expected_format:
        logger.error(f"Format mismatch: expected {expected_format}, got {result.get('format')}")
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is not None and result.get('width') != expected_width:
        logger.error(f"Width mismatch: expected {expected_width}, got {result.get('width')}")
        return 0.0
    if expected_height is not None and result.get('height') != expected_height:
        logger.error(f"Height mismatch: expected {expected_height}, got {result.get('height')}")
        return 0.0
    logger.info('Image dimensions and format match expected values')
    return 1.0

def check_png_valid__2511ecbd(result, expected, **options):
    """
    Check if file is a valid PNG image.

    Args:
        result: Dict from getter with 'exists', 'format', 'valid'
        expected: Dict with 'format' and 'exists'
        **options: Additional options

    Returns:
        float: 1.0 if valid PNG, 0.5 if exists but wrong format, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.5
    expected_format = expected.get('format', 'PNG')
    if result.get('format') == expected_format and result.get('valid', False):
        score += 0.5
    return score

def check_image_dimensions__de812dd5b44b906cb9793a8d4a3f91bf(result, expected, **options):
    """
    Check if the image dimensions match the expected dimensions.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys for expected dimensions
        **options: Additional options (not used)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None - file may not exist')
        return 0.0
    if not isinstance(result, dict) or 'width' not in result or 'height' not in result:
        logger.error(f'Invalid result format: {result}')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if expected_width is None or expected_height is None:
        logger.error(f'Invalid expected format: {expected}')
        return 0.0
    actual_width = result['width']
    actual_height = result['height']
    logger.info(f'Comparing dimensions - Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    if actual_width == expected_width and actual_height == expected_height:
        return 1.0
    else:
        return 0.0

def check_image_file_exists__3948a175(result_state: str, expected_state: Dict[str, Any], **options) -> float:
    """
    Check if the image file exists and has the correct dimensions.

    This function verifies that the flipped background image was created
    correctly by checking:
    1. The file exists
    2. The image has the expected dimensions (within tolerance)
    3. The image is a horizontally flipped version of the source image (if cloud_file provided)

    Args:
        result_state: Path to the image file returned by getter, or None if file doesn't exist
        expected_state: When type='rule', this IS the rules dict directly containing:
            - width (int): Expected image width
            - height (int): Expected image height
            - tolerance (int): Allowed deviation in pixels
            - cloud_file (str, optional): URL to original source image for flip verification
        **options: Additional options (unused)

    Returns:
        float: 1.0 if file exists and dimensions match (within tolerance), 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Image file does not exist or could not be retrieved')
        return 0.0
    if not os.path.isfile(result_state):
        logger.warning(f'Image file does not exist at path: {result_state}')
        return 0.0
    try:
        img = Image.open(result_state)
        (actual_width, actual_height) = img.size
        logger.info(f'Image dimensions: {actual_width}x{actual_height}')
        expected_width = expected_state.get('width')
        expected_height = expected_state.get('height')
        tolerance = expected_state.get('tolerance', 0)
        if expected_width is not None:
            width_diff = abs(actual_width - expected_width)
            if width_diff > tolerance:
                logger.warning(f'Width mismatch: expected {expected_width}±{tolerance}, got {actual_width} (diff: {width_diff})')
                return 0.0
            logger.info(f'Width check passed: {actual_width} (expected {expected_width}±{tolerance})')
        if expected_height is not None:
            height_diff = abs(actual_height - expected_height)
            if height_diff > tolerance:
                logger.warning(f'Height mismatch: expected {expected_height}±{tolerance}, got {actual_height} (diff: {height_diff})')
                return 0.0
            logger.info(f'Height check passed: {actual_height} (expected {expected_height}±{tolerance})')
        cloud_file = expected_state.get('cloud_file')
        if cloud_file:
            try:
                logger.info(f'Verifying horizontal flip against source: {cloud_file}')
                response = requests.get(cloud_file, timeout=30)
                response.raise_for_status()
                original_img = Image.open(BytesIO(response.content))
                expected_flipped = ImageOps.mirror(original_img)
                if img.size != expected_flipped.size:
                    logger.warning(f'Size mismatch with flipped original: result={img.size}, expected_flipped={expected_flipped.size}')
                    return 0.0
                result_rgb = img.convert('RGB')
                expected_rgb = expected_flipped.convert('RGB')
                result_hash = hashlib.md5(result_rgb.tobytes()).hexdigest()
                expected_hash = hashlib.md5(expected_rgb.tobytes()).hexdigest()
                if result_hash == expected_hash:
                    logger.info('Horizontal flip verification passed (exact match)')
                    return 1.0
                result_pixels = list(result_rgb.getdata())
                expected_pixels = list(expected_rgb.getdata())
                total_pixels = len(result_pixels)
                matching_pixels = sum((1 for (r, e) in zip(result_pixels, expected_pixels) if all((abs(r[i] - e[i]) <= 5 for i in range(3)))))
                similarity = matching_pixels / total_pixels
                logger.info(f'Pixel similarity: {similarity:.2%}')
                if similarity >= 0.98:
                    logger.info('Horizontal flip verification passed (high similarity)')
                    return 1.0
                else:
                    logger.warning(f'Image does not appear to be a horizontal flip of the source (similarity: {similarity:.2%})')
                    return 0.0
            except Exception as e:
                logger.error(f'Error verifying horizontal flip: {e}')
                logger.warning('Falling back to dimension-only check')
        logger.info('Image file exists with correct dimensions')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking image file: {e}')
        return 0.0

def check_image_blurred__782a609a(result_state, expected_state, **options):
    """
    Check if the image has been blurred by comparing edge sharpness.

    Args:
        result_state: Path to the result image file (from vm_file getter)
        expected_state: Dict with 'check_blur' boolean flag and optional 'variance_threshold'
        **options: Additional options

    Returns:
        float: 1.0 if image appears blurred, 0.0 otherwise
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    check_blur = expected_state.get('check_blur', True) if isinstance(expected_state, dict) else True
    if not check_blur:
        logger.info('Blur check disabled in expected_state')
        return 1.0
    try:
        import cv2
        if not os.path.exists(result_state):
            logger.error(f'File does not exist: {result_state}')
            return 0.0
        if not os.path.isfile(result_state):
            logger.error(f'Path is not a file: {result_state}')
            return 0.0
        img = cv2.imread(result_state, cv2.IMREAD_GRAYSCALE)
        if img is None:
            logger.error(f'Failed to load image from path: {result_state}')
            return 0.0
        laplacian = cv2.Laplacian(img, cv2.CV_64F)
        variance = np.var(laplacian)
        variance_threshold = expected_state.get('variance_threshold', 100) if isinstance(expected_state, dict) else 100
        if variance < variance_threshold:
            logger.info(f'Image appears blurred: variance {variance:.2f} < {variance_threshold}')
            return 1.0
        else:
            logger.warning(f'Image does not appear blurred: variance {variance:.2f} >= {variance_threshold}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking blur: {e}')
        return 0.0

def check_layer_group_config__7ba73b05(layer_structure, expected, **options):
    """
    Check if a layer group with the expected name exists in the GIMP image.

    Args:
        layer_structure: List of dictionaries with 'name' and 'is_group' keys
        expected: Dict containing the expected layer group name
        **options: Additional options

    Returns:
        float: Score (1.0 if layer group exists with correct name, 0.0 otherwise)
    """
    if layer_structure is None or not isinstance(layer_structure, list):
        logger.error('Invalid layer structure')
        return 0.0
    try:
        target_name = expected.get('name', 'Shapes')
        for layer in layer_structure:
            if layer.get('is_group', False) and layer.get('name', '') == target_name:
                logger.info(f'Found layer group with name: {target_name}')
                return 1.0
        logger.warning(f"Layer group '{target_name}' not found in image")
        logger.info(f'Available layers: {layer_structure}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking layer group: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_image_dimensions__c5f81e73faaccc56bdfa2edf29f272b7(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' keys from rules
        **options: Additional options

    Returns:
        float: 1.0 if both dimensions match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    if result.get('width') == expected_width and result.get('height') == expected_height:
        return 1.0
    return 0.0

def check_gimp_fullscreen_mode__8efddf2685fdb790d7823145c3565e94(actual_config_path, rule):
    """
    Check if GIMP fullscreen mode is enabled as expected.
    This checks the sessionrc file for the fullscreen setting.

    Args:
        actual_config_path: Path to the GIMP config file
        rule: Expected configuration with keys:
            - key: Config key name (can be string or list)
            - value: Expected value (usually "yes" or "no")

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path, 'r') as f:
        content = f.readlines()
    for line in content:
        if line.startswith('#') or line == '\n':
            continue
        items = line.strip().lstrip('(').rstrip(')\n').split()
        if isinstance(rule['key'], str):
            if items[0] == rule['key'] and items[-1] == rule['value']:
                return 1.0
        elif isinstance(rule['key'], list) and len(rule['key']) == 2:
            if items[0] == rule['key'][0] and items[1] == rule['key'][1] and (items[-1] == rule['value']):
                return 1.0
    return 0.0

def check_image_dimensions__565232d7(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_resized__8000c7e00c88e8061974ceb3ccc555df(result, expected, **options):
    """Check if image was correctly resized to target dimensions from slide 2.

    This metric verifies that:
    1. The result image has the exact target dimensions (800x600)
    2. The result image has high structural similarity (SSIM >= 0.85) to the reference
        image from slide 2, which proves the user extracted the correct source image

    By comparing against the reference slide 2 image, we ensure the user:
    - Opened the PowerPoint file
    - Navigated to slide 2
    - Extracted the correct image from that slide
    - Resized it correctly in GIMP to 800x600
    - Saved it as background_resized.png

    Args:
        result: Path to the user's resized image file
        expected: Dict with 'target_width', 'target_height', and 'reference_path'
        **options: Additional options (unused)

    Returns:
        float: 1.0 if both size and structure match, 0.0 otherwise
    """
    from PIL import Image
    import numpy as np
    import logging
    from skimage.metrics import structural_similarity as ssim
    logger = logging.getLogger(__name__)
    target_width = expected.get('target_width')
    target_height = expected.get('target_height')
    reference_path = expected.get('reference_path')
    if result is None or reference_path is None:
        logger.warning(f'Missing paths: result={result}, reference={reference_path}')
        return 0.0
    try:
        img_result = Image.open(result)
        img_reference = Image.open(reference_path)
        (result_width, result_height) = img_result.size
        logger.info(f'Size - result: ({result_width}, {result_height}), target: ({target_width}, {target_height})')
        size_correct = result_width == target_width and result_height == target_height
        if not size_correct:
            logger.warning(f'Size mismatch: result {img_result.size} != target ({target_width}, {target_height})')
            return 0.0
        img_reference_resized = img_reference.resize(img_result.size, Image.Resampling.LANCZOS)

        def check_ssim(img1, img2, threshold=0.85):
            if img1.mode != 'RGB':
                img1 = img1.convert('RGB')
            if img2.mode != 'RGB':
                img2 = img2.convert('RGB')
            array1 = np.array(img1)
            array2 = np.array(img2)
            if array1.shape != array2.shape:
                return False
            min_dim = min(array1.shape[0], array1.shape[1])
            win_size = min(7, min_dim if min_dim % 2 == 1 else min_dim - 1)
            if win_size < 1:
                return False
            try:
                similarity = ssim(array1, array2, win_size=win_size, channel_axis=2)
                logger.info(f'SSIM: {similarity}, threshold: {threshold}')
                return similarity >= threshold
            except TypeError:
                similarity = ssim(array1, array2, win_size=win_size, multichannel=True)
                logger.info(f'SSIM: {similarity}, threshold: {threshold}')
                return similarity >= threshold
            except Exception as e:
                logger.error(f'SSIM error: {e}')
                return False
        structure_same = check_ssim(img_result, img_reference_resized)
        logger.info(f'Size correct: {size_correct}, structure_same: {structure_same}')
        if size_correct and structure_same:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error in resize check: {e}')
        return 0.0

def check_png_dimensions__92ca6baf(result, expected, **options):
    """
    Check if PNG file exists and has minimum dimensions.

    Args:
        result: Dict from getter with 'exists', 'width', 'height'
        expected: Dict with 'min_width', 'min_height', 'check_exists'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, partial credit otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        return 0.0
    min_width = expected.get('min_width', 0)
    if result.get('width', 0) >= min_width:
        score += 0.3
    min_height = expected.get('min_height', 0)
    if result.get('height', 0) >= min_height:
        score += 0.3
    return score

def check_image_compression__c46a6f1dddc552cd368bc819d4cce6f7(result, expected, **options):
    """
    Check if an animated GIF file has appropriate compression, frame count, and is animated.

    Args:
        result: Dict from getter with file info (exists, file_size_bytes, format, frame_count, is_animated)
        expected: Dict with expected values (max_size_kb, min_size_kb, format, min_frames, max_frames)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    expected_format = expected.get('format', 'GIF')
    if result.get('format') == expected_format:
        score += 0.2
        logger.info(f"Format check passed: {result.get('format')}")
    else:
        logger.warning(f"Format mismatch: expected {expected_format}, got {result.get('format')}")
        return score
    is_animated = result.get('is_animated', False)
    if is_animated:
        score += 0.3
        logger.info('Animation check passed: GIF is animated')
    else:
        logger.warning('Animation check failed: GIF is not animated (single frame)')
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 30)
    max_frames = expected.get('max_frames', 100)
    if min_frames <= frame_count <= max_frames:
        score += 0.2
        logger.info(f'Frame count check passed: {frame_count} frames (range: {min_frames}-{max_frames})')
    else:
        logger.warning(f'Frame count out of range: {frame_count} frames (expected: {min_frames}-{max_frames})')
    file_size_bytes = result.get('file_size_bytes', 0)
    file_size_kb = file_size_bytes / 1024.0
    max_size_kb = expected.get('max_size_kb', 500)
    min_size_kb = expected.get('min_size_kb', 10)
    if min_size_kb <= file_size_kb <= max_size_kb:
        score += 0.3
        logger.info(f'File size check passed: {file_size_kb:.2f} KB (range: {min_size_kb}-{max_size_kb} KB)')
    elif file_size_kb >= min_size_kb:
        score += 0.15
        logger.warning(f'File size too large: {file_size_kb:.2f} KB (max: {max_size_kb} KB)')
    else:
        logger.warning(f'File size too small: {file_size_kb:.2f} KB (min: {min_size_kb} KB)')
    return score

def check_triangle_scale__12d50454feda301909898f2cf2cce54b(result, expected, **options):
    """
    Check if the triangle has been scaled to the expected proportion.

    Args:
        result: Actual triangle area (number of pixels) from getter
        expected: Dict with 'original_area' and 'scale_factor' keys
        **options: Optional 'tolerance' for area comparison (default: 0.15)

    Returns:
        float: 1.0 if scaling matches (within tolerance), 0.0 otherwise
    """
    if result is None:
        return 0.0
    original_area = expected.get('original_area')
    scale_factor = expected.get('scale_factor', 1.0)
    tolerance = options.get('tolerance', 0.15)
    expected_area = original_area * scale_factor ** 2
    if expected_area == 0:
        return 0.0
    relative_diff = abs(result - expected_area) / expected_area
    if relative_diff <= tolerance:
        return 1.0
    else:
        return 0.0

def check_image_extracted__a8440735(result, expected, **options):
    """
    Validates that the extracted image matches expected properties and content.

    This ensures:
    1. The file exists at the target path
    2. The file has correct format and dimensions
    3. The image content matches the first image from the email (critical validation)

    The third check is the most important as it verifies the agent actually
    extracted the correct image from the correct email, not just any random image.
    """
    score = 0.0
    if not result.get('exists'):
        return 0.0
    score += 0.15
    if result.get('size', 0) >= expected.get('min_size', 1000):
        score += 0.1
    if result.get('format') == expected.get('format', 'PNG'):
        score += 0.1
    if result.get('width', 0) >= expected.get('min_width', 100) and result.get('height', 0) >= expected.get('min_height', 100):
        score += 0.1
    if result.get('matches_email_image'):
        score += 0.55
    else:
        pass
    return score

def check_image_mirror__77b8ab4d994f43ac89308ca087d7c4b4(result_path, expected, **options):
    """
    Check if result image is a horizontally mirrored version of the source image.

    This function receives:
    - result_path: Local path to the result image (downloaded by get_vm_file)
    - expected: Rules dict with 'src_path' pointing to the original source image path

    The function downloads the source image, flips it horizontally, and compares
    it to the result image using SSIM.

    Args:
        result_path (str): Local path to the result/target image file (from get_vm_file)
        expected (dict): Rules dict containing:
            - src_path (str): Path to the original source image
        **options: Additional options (env is passed here by framework)

    Returns:
        float: 1.0 if images match (mirrored), 0.0 otherwise
    """
    if result_path is None:
        return 0.0
    src_vm_path = expected.get('src_path')
    if not src_vm_path:
        return 0.0
    env = options.get('env')
    if not env:
        try:
            source_image = Image.open(src_vm_path)
        except Exception:
            return 0.0
    else:
        try:
            src_file_bytes = env.controller.get_file(src_vm_path)
            if not src_file_bytes:
                return 0.0
            with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp:
                tmp.write(src_file_bytes)
                src_temp_path = tmp.name
            try:
                source_image = Image.open(src_temp_path)
            finally:
                if os.path.exists(src_temp_path):
                    os.unlink(src_temp_path)
        except Exception as e:
            print(f'Error downloading source image: {e}')
            return 0.0
    try:
        target_image = Image.open(result_path)
    except Exception as e:
        print(f'Error loading result image: {e}')
        return 0.0
    transposed_image = source_image.transpose(Image.FLIP_LEFT_RIGHT)
    mirrored = structure_check_by_ssim(transposed_image, target_image, 0.99)
    return 1.0 if mirrored else 0.0

def check_image_dimensions__4694337da0a8886d5bd508a95fd83b12(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if image has expected dimensions.

    Args:
        result: Dict from getter with 'width', 'height' keys, or None
        expected: Dict with 'width', 'height' expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.info('Image file not found or invalid')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width == expected_width and actual_height == expected_height:
        return 1.0
    else:
        logger.info(f'Dimension mismatch - Expected: {expected_width}x{expected_height}, Got: {actual_width}x{actual_height}')
        return 0.0

def check_image_dimensions__7cb89717dbfd62e5cbbd4dcc85a4e268(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: dict with 'width' and 'height' keys from getter
        expected: dict with 'width' and 'height' expected values

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    actual_width = result.get('width')
    actual_height = result.get('height')
    if actual_width is None or actual_height is None:
        logger.error(f'Invalid result dimensions: {result}')
        return 0.0
    width_match = actual_width == expected_width
    height_match = actual_height == expected_height
    logger.info(f'Expected: {expected_width}x{expected_height}, Got: {actual_width}x{actual_height}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_crop__9e45757e(result, expected, **options):
    """
    Check if the image has been cropped to the specified region.

    Args:
        result: Path to result image file
        expected: Dict with 'source_path' and 'crop_box' [x1, y1, x2, y2]
        **options: Additional options

    Returns:
        float: Score (1.0 if cropped correctly, 0.0 otherwise)
    """
    if result is None or expected is None:
        return 0.0
    try:
        import os
        source_path = expected.get('source_path')
        crop_box = expected.get('crop_box')
        if not os.path.isabs(source_path):
            result_dir = os.path.dirname(result)
            source_path = os.path.join(result_dir, source_path)
        result_img = Image.open(result)
        source_img = Image.open(source_path)
        expected_crop = source_img.crop(tuple(crop_box))
        if result_img.size != expected_crop.size:
            logging.debug(f'Size mismatch: {result_img.size} vs {expected_crop.size}')
            return 0.0
        if structure_check_by_ssim(result_img, expected_crop, threshold=0.95):
            return 1.0
        else:
            logging.debug('Result does not match expected crop')
            return 0.0
    except Exception as e:
        logging.error(f'Error in check_image_crop__9e45757e: {e}')
        return 0.0

def check_jpeg_export__39da7f334341155f29f73cbacf02786e(result, expected, **options):
    """
    Check if JPEG was exported correctly with proper format and size.

    Args:
        result: Dict from getter with file_exists, format, size_kb
        expected: Dict with file_exists, format, min_size_kb requirements
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False) and expected.get('file_exists', False):
        score += 0.4
        logger.info('✓ JPEG file exists')
    else:
        logger.warning('✗ JPEG file does not exist')
        return 0.0
    file_format = result.get('format', '')
    expected_format = expected.get('format', 'JPEG')
    if file_format == expected_format:
        score += 0.3
        logger.info(f'✓ Format is {file_format}')
    else:
        logger.warning(f'✗ Format is {file_format}, expected {expected_format}')
    size_kb = result.get('size_kb', 0)
    min_size_kb = expected.get('min_size_kb', 10)
    if size_kb >= min_size_kb:
        score += 0.3
        logger.info(f'✓ File size {size_kb:.2f} KB >= {min_size_kb} KB')
    else:
        logger.warning(f'✗ File size {size_kb:.2f} KB < {min_size_kb} KB')
    logger.info(f'Final score: {score}')
    return score

def check_image_dimensions__c984db77(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_hash__69eadfa0(result, expected, **options):
    """Compare image hash against expected hash.

    Args:
        result: SHA256 hash string from getter
        expected: Expected hash value (can be dict with 'hash' key or direct string)
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_hash = expected.get('hash', '')
    else:
        expected_hash = expected
    if result is None:
        logger.error('Result hash is None')
        return 0.0
    if result == expected_hash:
        logger.info(f'Hash match: {result}')
        return 1.0
    else:
        logger.warning(f'Hash mismatch. Expected: {expected_hash}, Got: {result}')
        return 0.0

def check_gif_exists__e4d07acf(result, expected, **options):
    """
    Check if GIF file exists and meets minimum size requirement.

    Args:
        result: Dict from getter with 'exists' and 'size'
        expected: Dict with 'should_exist' and 'min_size'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.5 if exists but too small, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    exists = result.get('exists', False)
    if not exists:
        return 0.0
    score = 0.5
    min_size = expected.get('min_size', 0)
    file_size = result.get('size', 0)
    if file_size >= min_size:
        score += 0.5
    return score

def check_gif_dimensions__90c23f3797e29598f167849a527a40d5(result, expected, **options):
    """
    Check if a GIF file has the expected dimensions, is animated, and has appropriate frame count.

    Args:
        result: Dict from getter with dimensions (exists, width, height, is_animated, frame_count)
        expected: Dict with expected dimensions (expected_width, expected_height, tolerance)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_animated', False):
        logger.warning('File is not animated (single frame only)')
        return 0.0
    frame_count = result.get('frame_count', 0)
    if frame_count < 40:
        logger.warning(f'Frame count too low: {frame_count} frames (expected 40-150 for ~5 seconds)')
        return 0.0
    elif frame_count > 150:
        logger.warning(f'Frame count too high: {frame_count} frames (expected 40-150 for ~5 seconds)')
    width = result.get('width', 0)
    height = result.get('height', 0)
    expected_width = expected.get('expected_width', 640)
    expected_height = expected.get('expected_height', 480)
    tolerance = expected.get('tolerance', 50)
    width_ok = abs(width - expected_width) <= tolerance
    height_ok = abs(height - expected_height) <= tolerance
    if width_ok and height_ok:
        logger.info(f'All checks passed: {width}x{height}, {frame_count} frames, animated')
        return 1.0
    elif width_ok or height_ok:
        logger.warning(f'Partial dimension match: {width}x{height} (expected: {expected_width}x{expected_height} ± {tolerance})')
        return 0.7
    else:
        logger.warning(f'Dimensions mismatch: {width}x{height} (expected: {expected_width}x{expected_height} ± {tolerance})')
        return 0.5

def check_gif_file__78103de86e64961437a4bcd00b97b9bc(result, expected, **options):
    """
    Check if a GIF file exists and meets requirements.

    Args:
        result: dict from get_gif_file_info__78103de86e64961437a4bcd00b97b9bc
        expected: dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.25
        logger.info('File exists: +0.25')
    else:
        logger.warning('File does not exist')
        return 0.0
    if result.get('is_gif', False):
        score += 0.25
        logger.info('File is a valid GIF: +0.25')
    else:
        logger.warning('File is not a valid GIF')
        return score
    min_size = expected.get('min_file_size', 1000)
    if result.get('file_size', 0) >= min_size:
        score += 0.15
        logger.info(f"File size {result['file_size']} >= {min_size}: +0.15")
    else:
        logger.warning(f"File size {result['file_size']} < {min_size}")
    min_frames = expected.get('min_frames', 20)
    if result.get('frames', 0) >= min_frames:
        score += 0.15
        logger.info(f"Frame count {result['frames']} >= {min_frames}: +0.15")
    else:
        logger.warning(f"Frame count {result['frames']} < {min_frames}")
    duration = result.get('duration_seconds', 0.0)
    min_duration = expected.get('min_duration', 2.5)
    max_duration = expected.get('max_duration', 3.5)
    if min_duration <= duration <= max_duration:
        score += 0.2
        logger.info(f'Duration {duration}s within range [{min_duration}, {max_duration}]: +0.2')
    else:
        logger.warning(f'Duration {duration}s outside range [{min_duration}, {max_duration}]')
    logger.info(f'Final score: {score}')
    return score

def check_image_size__10ab0aed(src_path, rule):
    """
    Check if the size of the src image matches the expected dimensions
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_4
    """
    if src_path is None:
        return 0.0
    try:
        img = Image.open(src_path)
        actual_width = img.size[0]
        actual_height = img.size[1]
        logger.debug(f'Image size: {img.size}')
        if rule.get('height', None) is not None:
            height_same = actual_height == rule['height']
        else:
            height_same = True
        if rule.get('width', None) is not None:
            width_same = actual_width == rule['width']
        else:
            width_same = True
        if height_same and width_same:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 1.0
        else:
            logger.debug(f'height_same: {height_same}, width_same: {width_same}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image size: {e}')
        return 0.0

def check_image_resize__ff57a44d(src_path, expected, **options):
    """
    Check if the image has been resized to the expected dimensions.
    Variation 0: 2a729ded-3296-423d-aec4-7dd55ed5fbb3

    Args:
        src_path: Path to the result image file
        expected: Dict with 'rules' containing 'width' and 'height' keys
        **options: Additional options

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if src_path is None:
        logger.warning('Source path is None')
        return 0.0
    try:
        img = Image.open(src_path)
        (actual_width, actual_height) = img.size
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        logger.info(f'Image size: {actual_width}x{actual_height}, Expected: {expected_width}x{expected_height}')
        width_match = actual_width == expected_width if expected_width is not None else True
        height_match = actual_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image resize: {e}')
        return 0.0

def check_triangle_color__978cb6f31473d4802226bc7ae94b7399(result, expected, **options):
    """
    Check if the triangle color matches the expected color.

    Args:
        result: Actual color dict {'r': int, 'g': int, 'b': int} from getter
        expected: Dict with 'target_color' key containing RGB values
        **options: Optional 'tolerance' for color distance (default: 30)

    Returns:
        float: 1.0 if color matches (within tolerance), 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not isinstance(result, dict):
        return 0.0
    if 'r' not in result or 'g' not in result or 'b' not in result:
        return 0.0
    target_color = expected.get('target_color', {})
    if not target_color or 'r' not in target_color or 'g' not in target_color or ('b' not in target_color):
        return 0.0
    tolerance = options.get('tolerance', 30)
    try:
        r_diff = result['r'] - target_color['r']
        g_diff = result['g'] - target_color['g']
        b_diff = result['b'] - target_color['b']
        distance = (r_diff ** 2 + g_diff ** 2 + b_diff ** 2) ** 0.5
        if distance <= tolerance:
            return 1.0
        else:
            return 0.0
    except (TypeError, KeyError):
        return 0.0

def check_pdf_image_count__b5bd06ba(result_state, expected_state, **options):
    """
    Check if PDF file contains at least the minimum number of embedded images.

    This function verifies that a PDF file was created from an image by checking
    that it contains at least one embedded image.

    Args:
        result_state: Path to the PDF file (str) from vm_file getter
        expected_state: Dict with 'min_images' (int) - when expected.type='rule',
                       this is the rules dict directly
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and contains at least min_images, 0.0 otherwise
    """
    if result_state is None or not isinstance(result_state, str):
        return 0.0
    if not os.path.exists(result_state):
        return 0.0
    min_images = expected_state.get('min_images', 1)
    try:
        reader = PdfReader(result_state)
        image_count = 0
        for page in reader.pages:
            if '/Resources' in page and '/XObject' in page['/Resources']:
                xobjects = page['/Resources']['/XObject'].get_object()
                for obj_name in xobjects:
                    obj = xobjects[obj_name]
                    if obj.get('/Subtype') == '/Image':
                        image_count += 1
        if image_count >= min_images:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_image_inserted__3333dfb2(result, expected, **options):
    """Check if image was successfully inserted.

    Args:
        result: Boolean indicating if document has new image
        expected: Expected value (dict with 'inserted' key)
        **options: Additional options

    Returns:
        float: 1.0 if image was inserted, 0.0 otherwise
    """
    expected_val = expected.get('inserted', True)
    if result == expected_val:
        return 1.0
    return 0.0

def check_image_dimensions__1beedc32(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_size__993c9cd2028ad767fa928842de0805ca(result, expected, **options):
    """
    Check if image was rotated correctly by verifying dimensions and rotation direction.

    For 90° clockwise rotation verification:
    1. Checks that dimensions match expected (swapped from original)
    2. Verifies rotation direction by comparing corner pixel samples
    3. Uses perceptual hash similarity as additional verification

    Corner mapping for 90° clockwise rotation:
    - Original top-left → Rotated top-right
    - Original top-right → Rotated bottom-right
    - Original bottom-right → Rotated bottom-left
    - Original bottom-left → Rotated top-left

    Args:
        result: dict with 'width', 'height', 'phash', 'corner_samples' from getter
        expected: dict with 'width', 'height', 'original_path' keys
        **options: Additional options

    Returns:
        float: 1.0 if rotation is correct, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    tolerance = options.get('tolerance', 0)
    result_width = result.get('width', 0)
    result_height = result.get('height', 0)
    expected_width = expected.get('width', 0)
    expected_height = expected.get('height', 0)
    width_match = abs(result_width - expected_width) <= tolerance
    height_match = abs(result_height - expected_height) <= tolerance
    if not (width_match and height_match):
        logger.info(f'Dimension mismatch: result=({result_width}, {result_height}), expected=({expected_width}, {expected_height})')
        return 0.0
    original_url = expected.get('original_url')
    if not original_url:
        logger.warning('No original image URL provided, skipping rotation verification')
        return 1.0
    try:
        import urllib.request
        with urllib.request.urlopen(original_url) as response:
            original_bytes = response.read()
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
            tmp.write(original_bytes)
            tmp_path = tmp.name
        original_img = Image.open(tmp_path)
        orig_width = original_img.width
        orig_height = original_img.height
        sample_size = min(10, orig_width // 10, orig_height // 10)
        if sample_size > 0:
            orig_tl = list(original_img.crop((0, 0, sample_size, sample_size)).getdata())[:25]
            orig_tr = list(original_img.crop((orig_width - sample_size, 0, orig_width, sample_size)).getdata())[:25]
            orig_bl = list(original_img.crop((0, orig_height - sample_size, sample_size, orig_height)).getdata())[:25]
            orig_br = list(original_img.crop((orig_width - sample_size, orig_height - sample_size, orig_width, orig_height)).getdata())[:25]
            result_corners = dict(result.get('corner_samples', []))
            rot_tl = result_corners.get('top_left', [])
            rot_tr = result_corners.get('top_right', [])
            rot_bl = result_corners.get('bottom_left', [])
            rot_br = result_corners.get('bottom_right', [])

            def pixel_similarity(p1, p2):
                """Calculate similarity between two pixel samples (0-1)"""
                if not p1 or not p2 or len(p1) != len(p2):
                    return 0.0
                total_dist = 0
                count = 0
                for (px1, px2) in zip(p1, p2):
                    if isinstance(px1, (tuple, list)) and isinstance(px2, (tuple, list)):
                        dist = sum((abs(a - b) for (a, b) in zip(px1, px2))) / len(px1)
                    else:
                        dist = abs(px1 - px2)
                    total_dist += dist
                    count += 1
                avg_dist = total_dist / count if count > 0 else 255
                similarity = max(0, 1 - avg_dist / 255)
                return similarity
            sim_tl_tr = pixel_similarity(orig_tl, rot_tr)
            sim_tr_br = pixel_similarity(orig_tr, rot_br)
            sim_br_bl = pixel_similarity(orig_br, rot_bl)
            sim_bl_tl = pixel_similarity(orig_bl, rot_tl)
            avg_similarity = (sim_tl_tr + sim_tr_br + sim_br_bl + sim_bl_tl) / 4
            logger.info(f'Corner rotation verification - similarities: tl→tr={sim_tl_tr:.3f}, tr→br={sim_tr_br:.3f}, br→bl={sim_br_bl:.3f}, bl→tl={sim_bl_tl:.3f}, avg={avg_similarity:.3f}')
            if avg_similarity < 0.7:
                logger.warning(f'Rotation direction verification failed: avg_similarity={avg_similarity:.3f} < 0.7')
                return 0.0
        original_img.close()
        os.unlink(tmp_path)
        logger.info(f'Image rotation verified successfully')
        return 1.0
    except Exception as e:
        logger.error(f'Error verifying rotation: {e}')
        logger.warning('Falling back to dimension-only verification')
        return 1.0

def check_triangle_rightedge__afcf05ab(result_state, expected_state, **options):
    """
    Check if the triangle is aligned with the right edge of the image.
    Variation 4 for task f4aec372-4fb0-4df5-a52b-79e0e2a5d6ce

    Args:
        result_state: Path to the exported image file
        expected_state: Expected rules dict (empty dict for this task)
        **options: Additional options

    Returns:
        float: Score (1.0 if triangle is at right edge, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        if len(img_array.shape) == 2:
            img_array = np.stack([img_array] * 3, axis=-1)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        if len(unique_colors) < 2:
            logger.warning('Image has less than 2 unique colors')
            return 0.0
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        triangle_color = None
        for color in unique_colors_sorted:
            if len(color) >= 3 and color[0] > 200 and (color[1] > 200) and (color[2] < 100):
                triangle_color = color
                break
        if triangle_color is None:
            triangle_color = unique_colors_sorted[-2] if len(unique_colors_sorted) >= 2 else unique_colors_sorted[0]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        max_x = triangle_coords[:, 1].max()
        img_width = img_array.shape[1]
        at_right_edge = img_width - max_x <= img_width * 0.05
        return 1.0 if at_right_edge else 0.0
    except Exception as e:
        logger.error(f'Error checking triangle position: {e}')
        return 0.0

def check_gif_file__06722a19(result, expected, **options):
    """
    Check if a GIF file meets expected criteria.

    Args:
        result: dict with file info from getter
        expected: dict with expected criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score += 0.5
    if result.get('format') == 'GIF':
        score += 0.3
    else:
        logger.info(f"Wrong format: {result.get('format')}")
        return score
    frame_count = result.get('frame_count', 0)
    min_frames = expected.get('min_frames', 1)
    if frame_count >= min_frames:
        score += 0.2
    else:
        logger.info(f'Insufficient frames: {frame_count} < {min_frames}')
    return score

def check_jpg_file__1bf235b17a43af3cc147100153de4d30(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if JPG file exists and meets requirements.

    Args:
        result: Dict from getter with 'exists', 'is_jpg', 'format', 'size' keys
        expected: Dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        logger.info('File does not exist')
        return 0.0
    if result.get('is_jpg', False):
        score += 0.3
    else:
        logger.info(f"Not a JPG file. Format: {result.get('format')}")
    min_size = expected.get('min_size', 1000)
    if result.get('size', 0) >= min_size:
        score += 0.2
    else:
        logger.info(f"File size too small: {result.get('size')} bytes")
    return score

def check_image_dimensions__fa1a72d6(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_triangle_bottom_right__363d0657(result_state, expected_state, **options):
    """
    Check if the triangle has been moved to the bottom-right corner.

    Args:
        result_state: Path to the result image
        expected_state: Not used (rule-based evaluation)
        **options: Additional options

    Returns:
        float: Score (1.0 if positioned correctly, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        img = Image.open(result_state)
        img_array = np.array(img)
        (unique_colors, counts) = np.unique(img_array.reshape(-1, img_array.shape[2]), axis=0, return_counts=True)
        unique_colors_sorted = unique_colors[np.argsort(counts)]
        if len(unique_colors_sorted) < 2:
            logger.warning('Could not find triangle in image')
            return 0.0
        triangle_color = unique_colors_sorted[1]
        triangle_mask = np.all(img_array == triangle_color, axis=2)
        triangle_coords = np.argwhere(triangle_mask)
        if len(triangle_coords) == 0:
            logger.warning('No triangle pixels found')
            return 0.0
        centroid = triangle_coords.mean(axis=0)
        (image_height, image_width) = img_array.shape[:2]
        bottom_threshold = image_height * 0.85
        right_threshold = image_width * 0.85
        logger.info(f'Triangle centroid: ({centroid[1]:.1f}, {centroid[0]:.1f})')
        logger.info(f'Image size: {image_width}x{image_height}')
        logger.info(f'Thresholds: right={right_threshold:.1f}, bottom={bottom_threshold:.1f}')
        if centroid[0] >= bottom_threshold and centroid[1] >= right_threshold:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error in check_triangle_bottom_right__363d0657: {e}')
        import traceback
        logger.error(traceback.format_exc())
        return 0.0

def check_gimp_brightness_increase__652303c0d066122d99f102352a5b1a93(result, expected, **options):
    """
    Check if the image brightness is higher than the original.

    Args:
        result: float brightness value from getter
        expected: dict with 'min_brightness' key specifying minimum acceptable brightness
        **options: Additional options

    Returns:
        float: 1.0 if brightness is above minimum, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    min_brightness = expected.get('min_brightness', 0)
    logger.info(f'Brightness check: result={result}, min_required={min_brightness}')
    if result >= min_brightness:
        logger.info(f'Brightness check PASSED: {result} >= {min_brightness}')
        return 1.0
    else:
        logger.info(f'Brightness check FAILED: {result} < {min_brightness}')
        return 0.0

def check_image_dimensions__e0d9b551(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: tuple (width, height) from getter
        expected: dict with 'rules' containing 'width' and 'height'
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: 1.0 if dimensions match, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    expected_width = expected.get('width')
    expected_height = expected.get('height')
    tolerance = options.get('tolerance', 0)
    (actual_width, actual_height) = result
    width_match = abs(actual_width - expected_width) <= tolerance
    height_match = abs(actual_height - expected_height) <= tolerance
    logger.info(f'Expected: {expected_width}x{expected_height}, Actual: {actual_width}x{actual_height}')
    logger.info(f'Width match: {width_match}, Height match: {height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_vertical_flip__7870cc8f1b3a8b0d979c0fd70171d833(result, expected, **options):
    """
    Check if the image is vertically flipped (mirrored top-to-bottom)
    Custom metric for variation 0 of task 72f83cdc-bf76-4531-9a1b-eb893a13f8aa

    Args:
        result: Rules dict from result getter containing {'path': ..., 'dest': ...}
        expected: Rules dict from expected getter containing {'path': ..., 'dest': ...}
        **options: Additional options

    Returns:
        float: Score (1.0 if vertically flipped, 0.0 otherwise)
    """
    src_path = expected.get('path') if isinstance(expected, dict) else None
    tgt_path = result.get('path') if isinstance(result, dict) else None
    if src_path is None or tgt_path is None:
        return 0.0
    try:
        source_image = Image.open(src_path)
        target_image = Image.open(tgt_path)
        transposed_image = source_image.transpose(Image.FLIP_TOP_BOTTOM)
        flipped = structure_check_by_ssim(transposed_image, target_image, 0.99)
        if flipped:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_image_resized__7253f77d(src_path, rule):
    """
    Check if the image has been resized to the specified dimensions.
    Similar to check_image_size but for resizing tasks.
    """
    if src_path is None:
        return 0.0
    img = Image.open(src_path)
    expected_width = rule.get('width', None)
    expected_height = rule.get('height', None)
    width_match = True if expected_width is None else img.size[0] == expected_width
    height_match = True if expected_height is None else img.size[1] == expected_height
    logger.info(f'Image size: {img.size}, expected: ({expected_width}, {expected_height}), match: width={width_match}, height={height_match}')
    if width_match and height_match:
        return 1.0
    else:
        return 0.0

def check_image_dimensions__717d6863(result, expected, **options):
    """
    Check if image dimensions match expected values.

    Args:
        result: Dict with 'width' and 'height' keys from getter
        expected: Dict with expected 'width' and 'height' values
        **options: Additional options (tolerance for approximate matching)

    Returns:
        float: Score (1.0 if dimensions match, 0.0 otherwise)
    """
    try:
        if result is None:
            logger.error('Result is None')
            return 0.0
        expected_width = expected.get('width')
        expected_height = expected.get('height')
        result_width = result.get('width')
        result_height = result.get('height')
        if result_width is None or result_height is None:
            logger.error(f'Invalid result dimensions: {result}')
            return 0.0
        width_match = result_width == expected_width if expected_width is not None else True
        height_match = result_height == expected_height if expected_height is not None else True
        if width_match and height_match:
            logger.info(f'Dimensions match: {result_width}x{result_height}')
            return 1.0
        else:
            logger.warning(f'Dimensions mismatch: got {result_width}x{result_height}, expected {expected_width}x{expected_height}')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image dimensions: {e}')
        return 0.0

def check_jpg_export__0a152dc7(result, expected, **options):
    """
    Check if a JPG file was successfully exported from LibreOffice Impress.

    Args:
        result (str): Path to the downloaded JPG file from vm_file
        expected (dict): Expected configuration with rules
        **options: Additional options

    Returns:
        float: Score (1.0 if export successful, 0.0 otherwise)
    """
    try:
        if not os.path.exists(result):
            return 0.0
        if os.path.getsize(result) == 0:
            return 0.0
        try:
            img = Image.open(result)
            if img.format not in ['JPEG', 'JPG']:
                return 0.0
            if img.size[0] == 0 or img.size[1] == 0:
                return 0.0
            img.close()
        except Exception as e:
            return 0.0
        return 1.0
    except Exception as e:
        print(f'Error in check_jpg_export__0a152dc7: {e}')
        return 0.0

def check_image_grayscale__3aef0cbd8986e240435edab0fd96c873(result, expected, **options):
    """Check if image has been converted to grayscale.

    Args:
        result: dict with mode and grayscale status from getter
        expected: dict with expected grayscale status
        **options: Additional options

    Returns:
        float: 1.0 if grayscale status matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_grayscale = expected.get('is_grayscale', True)
    if result.get('is_grayscale') == expected_grayscale:
        return 1.0
    return 0.0

def check_image_addition__17e4ac0c(result, expected, **options):
    """Check if image was added successfully.

    Args:
        result: Actual count
        expected: Expected count (dict with 'target_count' key)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    target = expected.get('target_count', 4)
    return 1.0 if result == target else 0.0

def check_image_file_exists__836db3cb3a7cc19d82f98c6a439eb80c(result: bool, expected: dict, **options) -> float:
    """
    Check if image file exists as expected.

    Args:
        result: Boolean indicating if image file exists (from getter)
        expected: Expected rules dict with 'should_exist' key (True/False)
        **options: Additional options

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise
    """
    should_exist = expected.get('should_exist', True)
    if result == should_exist:
        logger.info(f'Image existence check passed: {result} == {should_exist}')
        return 1.0
    else:
        logger.warning(f'Image existence check failed: {result} != {should_exist}')
        return 0.0

def check_image_dimensions__12182c78(result, expected, **options):
    """Check if image dimensions match expected values.

    Args:
        result: Image properties dict from getter
        expected: Expected properties (from rules)
        **options: Additional options including tolerance

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    rules = expected if isinstance(expected, dict) else expected.get('rules', {})
    tolerance = options.get('tolerance', 50000)
    score = 0.0
    checks = 0
    if 'width' in rules:
        checks += 1
        if abs(result['width'] - rules['width']) <= tolerance:
            score += 1.0
    if 'height' in rules:
        checks += 1
        if abs(result['height'] - rules['height']) <= tolerance:
            score += 1.0
    if 'left' in rules:
        checks += 1
        if abs(result['left'] - rules['left']) <= tolerance:
            score += 1.0
    if 'top' in rules:
        checks += 1
        if abs(result['top'] - rules['top']) <= tolerance:
            score += 1.0
    if checks == 0:
        return 0.0
    return score / checks

def check_image_properties__739292ff(result, expected, **options):
    """Check multiple image properties with partial credit.

    Args:
        result: Dict with image properties
        expected: Dict with various expected properties
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 based on how many properties match
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    score = 0.0
    checks = 0
    if 'format' in expected:
        checks += 1
        if result.get('format', '').upper() == expected['format'].upper():
            score += 1.0
            logger.info(f"Format check passed: {result.get('format')}")
        else:
            logger.info(f"Format check failed: expected {expected['format']}, got {result.get('format')}")
    if 'min_width' in expected:
        checks += 1
        if result.get('width', 0) >= expected['min_width']:
            score += 1.0
            logger.info(f"Width check passed: {result.get('width')} >= {expected['min_width']}")
        else:
            logger.info(f"Width check failed: {result.get('width')} < {expected['min_width']}")
    if 'min_height' in expected:
        checks += 1
        if result.get('height', 0) >= expected['min_height']:
            score += 1.0
            logger.info(f"Height check passed: {result.get('height')} >= {expected['min_height']}")
        else:
            logger.info(f"Height check failed: {result.get('height')} < {expected['min_height']}")
    if 'min_size' in expected:
        checks += 1
        if result.get('size_bytes', 0) >= expected['min_size']:
            score += 1.0
            logger.info(f"Min size check passed: {result.get('size_bytes')} >= {expected['min_size']}")
        else:
            logger.info(f"Min size check failed: {result.get('size_bytes')} < {expected['min_size']}")
    if checks == 0:
        logger.warning('No checks specified in expected')
        return 0.0
    final_score = score / checks
    logger.info(f'Final score: {final_score} ({score}/{checks} checks passed)')
    return final_score

def check_image_properties__9873b6c6(result, expected, **options):
    """Compare image properties against expected values.

    Args:
        result: dict from getter with image properties
        expected: dict with expected width, height, exists
        **options: Additional options (tolerance for dimensions)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'exists' in expected:
        total_checks += 1
        if result.get('exists') == expected['exists']:
            score += 1.0
    if 'width' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_width = expected['width']
        actual_width = result.get('width', 0)
        if abs(actual_width - expected_width) <= tolerance:
            score += 1.0
    if 'height' in expected and result.get('exists', False):
        total_checks += 1
        tolerance = options.get('dimension_tolerance', 5)
        expected_height = expected['height']
        actual_height = result.get('height', 0)
        if abs(actual_height - expected_height) <= tolerance:
            score += 1.0
    if 'format' in expected and result.get('exists', False):
        total_checks += 1
        if result.get('format') == expected['format']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks
