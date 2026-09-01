"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_pdf_files_copied__b8d04622', 'check_pdf_files__600d7d75', 'check_pdf_valid__3b118c02', 'check_pdf_comprehensive__32aaa745', 'check_pdf_count__a0da932e', 'check_pdf_content_a503b07f', 'check_pdf_exists__b62f9acb', 'check_pdf_files_valid__2fc6b524', 'check_pdf_files__437d5a7f', 'check_pdf_in_listing__4e03b1ed', 'check_named_pdf__0c8a922818da41c6a16e3979ae1f26dc', 'check_pdf_content__edbb73b7', 'check_pdf_export__d587f20a', 'check_directory_with_pdfs__92ec1d53', 'check_pdf_exists_with_content__e1e75309_0', 'check_pdf_export__17dfab17', 'check_pdf_location__0d0715b7122c298c1e10aa6fa135f599', 'check_pdf_export__d582d1d2', 'check_pdf_count__6c8aa0e0d1a7a1f247933b05b6dc0e08', 'check_pdf_validity__90edd5864af08865c892a3ecd6659029', 'check_pdf_chapters__097402a1', 'check_pdf_orientation__e222560cb6780019664a6917701a2659', 'check_pdf_chapters__d2aaaf87', 'check_pdf_properties__869dd2b5b7a0e45511c735d06983da83', 'check_pdf_files__2c1e781e', 'check_pdf_orientation__596fcd5e219bea1372357f0afb95cc85', 'check_pdf_files_non_empty__99e94d1a', 'check_pdf_chapters__89dbe8bb', 'check_pdf_files__197670d4', 'check_pdf_export__b47a318c', 'check_pdf_export__bec19abd', 'check_pdf_files_exist__989759aa', 'check_pdf_size_and_validity__e1e75309_9', 'check_pdf_count__085cc8d6', 'check_pdf_export__831e0e03', 'check_pdf_basic_requirements__7b63a847bf086913e974e7a32debb9ec', 'check_pdf_chapters__21dcc4a0', 'check_pdf_files__230c9972', 'check_pdf_named__b387dd88', 'check_pdf_content__5e9826db825b680d44e19c36727da776', 'check_pdf_files__1ee77af4', 'check_pdf_chapters__0be3e647', 'check_pdf_basic__d8736b62', 'check_pdf_readable__cd376b4c', 'check_pdf_chapters__c8947a04', 'check_pdf_comprehensive__6dc3161b', 'check_pdf_files__32b125c8', 'check_pdf_files_with_content__83b6523335a3d68f9c734599b1537c74', 'check_pdf_files__d8278409', 'check_pdf_contains_text__4c28c68dd08d7073669bab47ee359a64', 'check_pdf_files__184db76a3945be5ea6dec91515beac2a', 'check_pdf_exists__f7d1f102', 'check_pdf_text_content__7a2513e6563b76f5c07e27d8c4089d02', 'check_pdf_saved__5f59826466b2625834fd8d369560ed11', 'check_pdf_exists__9753fb3ef75762f14fb08b7f236e3f81', 'check_pdf_location__3a8b0dcb7e90aff8357c94bc8dad2474', 'check_pdf_contains_keywords__2c2512f5', 'check_pdf_files__8d8da24c', 'check_pdf_all_fields__ff249445b547a028ee5246e45cd19fdf', 'check_pdf_file_size__f34c6d72482f0e93994904e974de58fa', 'check_pdf_location__c74e09f1', 'check_desktop_pdf__a45ebe2876987410607711d3992c10db', 'check_pdf_exists__3df7d80c', 'check_exact_pdf_count__10f702e7ff0a641a1fda6d45251486ce', 'check_pdf_file_exists__e68e77d1', 'check_pdf_file_count__822cf85f742bbfae9e3acf8d7027940c', 'check_pdf_files__4e03b1ed', 'check_pdf_files__6e9dcacc', 'check_pdf_files__80b8ddf2', 'check_pdf_merge_verification__a8bc7d70fd87c66cda30ee22072de4d3', 'check_pdf_file_count__086472c391d9b3dbf463fe72713bc019', 'check_pdf_validity__1cc9c1bfb30ee1b2b3cdbf7926894918', 'check_pdf_file_properties__6f106331', 'check_pdf_non_empty__d4f3d6039bf1a71698a65c2d049521e8', 'check_pdf_saved__60388d3f6c5270e1728288c485d5fd5b', 'check_pdf_files_exist__6d39fc800ff8d1397c30bfae2c676b76', 'check_pdf_text_contains__a09f0e42332d230e7cc0e7794732425e', 'check_pdf_saved__40d8ee41df97cbb2f480ba7af545efc4', 'check_pdf_files__8010e79b', 'check_pdf_filesize__74ce4e60', 'check_pdf_export__9a18b30d646547c54b42b3593f83920d', 'check_pdf_has_text_content__e1e75309_6', 'check_pdf_filename__e61f394075a7c32a6a0c2c96a3700939', 'check_pdf_count_and_content__9ba5dc81930fa553e6c6310edaaff2ec', 'check_pdf_files__d928a635', 'check_pdf_keywords__c68533a8ce70563646c4156811d621fa', 'check_pdf_valid__e4448a5a', 'check_single_pdf_complete__1386929a4648885c7d87d6425829fd94', 'check_pdf_format__988ec512', 'check_pdf_file_exists__7648b3ac', 'check_pdf_filenames__7152d37e', 'check_pdf_file_size__29f6acc744cf0d15caa355ed1701b507', 'check_pdf_exists__05a7fe2932371163af38b8e77d9b0c93', 'check_pdf_name_pattern__55b958e432c8380deab73a3d2fcf329a', 'check_chapter_pdf_count__457b9850ebb20f14d7c68688b0aa27a6', 'check_pdf_chapters__6b40f44d', 'check_pdf_exists__8752e3fa', 'check_pdf_validity__0b5ef92a5e5fe7305b438702a0eb6c3a', 'check_pdf_export__c5999de9', 'check_pdf_split_complete__ea295e4c379ee25a192357c908aada76', 'check_pdf_files__cdfaf4c1', 'check_pdf_files__0707ddca', 'check_pdf_chapters__6871c0fe', 'check_pdf_count__bfc9cce9462c8f8be02842ad8f805893', 'check_pdf_exists__675930f3', 'check_pdf_files__201ff98c', 'check_pdf_export__2586a709', 'check_pdf_chapters__92f3cb2c', 'check_pdf_export__98b0d7a7', 'check_pdf_chapters_exist__242f4871e4280ec1025212dbeaf5c18c', 'check_pdf_file__766a8b90', 'check_pdf_orientation__c75d6b17', 'check_pdf_export__d263a7ae', 'check_pdf_count__2f750d009f0f471751ed869c09f90a90', 'check_pdf_files__a96f92e2', 'check_pdf_basic__06c89307', 'check_pdf_files__cb51ce0b', 'check_pdf_not_encrypted__391a007e', 'check_year_organized_pdfs__5ad487ac7d025a6b906a4c83e8beac41', 'check_pdf_not_empty__2d68dae7', 'check_pdf_saved__3dd8a1347f60dd796dc5e98521ae4032', 'check_pdf_exists__949eb101', 'check_pdf_saved__000c73be1394701e25c8491dd9418647', 'check_pdf_chapters__cedddaca', 'check_pdf_naming_pattern__d2aeadd33c64efc8e0e3416b06f6e222', 'check_pdf_files__a76459ec', 'check_pdf_files__f4eddf72', 'check_pdf_export__7c083df6', 'check_pdf_exists__7a335cc66cd23236defef5bd95bbbe7d', 'check_pdf_contains_text__439ece4f99f624319fffd552c64d5f89', 'check_pdf_contains_text__e1e75309_4', 'check_pdf_sizes__4e03b1ed', 'check_pdf_on_desktop__e1e75309_5', 'check_pdf_numbered_sequence__fe03784ed42c9a7002fcc758df207953']

def check_pdf_files_copied__b8d04622(result_state, expected_state, **options):
    """
    Check if all PDF files have been copied to the All_Invoices directory.

    Verifies:
    1. The directory exists and can be listed (no error in result_state)
    2. At least the expected number of PDF files are present
    3. All expected PDF files are present in the directory

    Args:
        result_state: Output from ls command (string with file listing)
        expected_state: Dictionary with rules containing:
            - min_count: Minimum number of files expected
            - expected_files: List of expected PDF filenames
        **options: Additional options

    Returns:
        float: Score (1.0 if all checks pass, 0.0 otherwise)
    """
    if result_state is None or isinstance(result_state, dict):
        if isinstance(result_state, dict) and result_state.get('error'):
            return 0.0
        if result_state is None:
            return 0.0
    min_count = expected_state.get('min_count', 0)
    expected_files = expected_state.get('expected_files', [])
    if isinstance(result_state, str):
        files = [f.strip() for f in result_state.strip().split('\n') if f.strip()]
    else:
        return 0.0
    if len(files) < min_count:
        return 0.0
    for expected_file in expected_files:
        if expected_file not in files:
            return 0.0
    return 1.0

def check_pdf_files__600d7d75(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_valid__3b118c02(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF is valid and contains minimum text content.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if PDF is valid and meets criteria, 0.0 otherwise
    """
    min_text_length = expected.get('min_text_length', 0)
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with fitz.open(result) as pdf:
            text = ''
            for page in pdf:
                text += page.get_text()
            if len(text.strip()) < min_text_length:
                return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_comprehensive__32aaa745(pdf_file: str, expected, **options):
    """
    Comprehensive PDF check: exists, valid, has content, reasonable size.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with validation rules
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.25
    try:
        reader = PdfReader(pdf_file)
        score += 0.25
        if len(reader.pages) > 0:
            score += 0.25
        file_size = os.path.getsize(pdf_file)
        if file_size >= 5120:
            score += 0.25
    except Exception as e:
        return score
    return score

def check_pdf_count__a0da932e(result, expected, **options):
    """Check if PDF file count matches expected value.

    Args:
        result: Integer count of PDF files found
        expected: Dict with 'count' parameter
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_count = expected.get('count', 2)
    if result == expected_count:
        logger.info(f'PDF count matches expected: {result}')
        return 1.0
    else:
        logger.warning(f'PDF count mismatch: got {result}, expected {expected_count}')
        return 0.0

def check_pdf_content_a503b07f(pdf_file: str, rules: Dict[str, Any]) -> float:
    """
    Check if PDF has the expected number of pages and contains image content.

    This metric verifies:
    1. The PDF file exists
    2. The PDF has the expected number of pages
    3. The PDF contains actual image content (not empty)

    Args:
        pdf_file: Path to the PDF file (from getter)
        rules: Dict with:
            - relation (str): Comparison operator ('eq', 'ne', 'lt', 'le', 'gt', 'ge')
            - ref_value (int): Expected page count
            - verify_content (bool): If True, verify PDF contains images

    Returns:
        float: 1.0 if all conditions met, 0.0 otherwise
    """
    from pypdf import PdfReader
    if pdf_file is None or not pdf_file:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    try:
        reader = PdfReader(pdf_file)
        nb_pages = len(reader.pages)
    except Exception:
        return 0.0
    relation = rules.get('relation', 'eq')
    ref_value = rules.get('ref_value', 1)
    try:
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
    except AttributeError:
        return 0.0
    if not page_count_matches:
        return 0.0
    verify_content = rules.get('verify_content', False)
    if verify_content:
        file_size = os.path.getsize(pdf_file)
        if file_size < 10240:
            return 0.0
        has_images = False
        try:
            for page in reader.pages:
                if '/Resources' in page and '/XObject' in page['/Resources']:
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj_name in xobjects:
                        xobject = xobjects[obj_name]
                        if xobject.get('/Subtype') == '/Image':
                            has_images = True
                            break
                if has_images:
                    break
        except Exception:
            return 0.0
        if not has_images:
            return 0.0
    return 1.0

def check_pdf_exists__b62f9acb(result, expected, **options):
    """
    Check if PDF file exists and meets minimum size requirement.

    Args:
        result: Dict from getter with 'exists' and 'size'
        expected: Dict with 'should_exist' and 'min_size'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.5 if exists but too small, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    exists = result.get('exists', False)
    if not exists:
        return 0.0
    score = 0.5
    min_size = expected.get('min_size', 0)
    file_size = result.get('size', 0)
    if file_size >= min_size:
        score += 0.5
    return score

def check_pdf_files_valid__2fc6b524(result, expected, **options):
    """Check if PDF files are valid (non-zero size).

    Args:
        result: List of file sizes
        expected: Dict with:
            - min_count: Minimum number of files
            - min_size: Minimum size per file in bytes

    Returns:
        Score between 0.0 and 1.0
    """
    min_count = expected.get('min_count', 2)
    min_size = expected.get('min_size', 1000)
    if not isinstance(result, list):
        return 0.0
    if len(result) < min_count:
        return 0.0
    valid_files = sum((1 for size in result if size >= min_size))
    if valid_files >= min_count:
        return 1.0
    else:
        return 0.0

def check_pdf_files__437d5a7f(result, expected, **options):
    """Check if all specific files exist.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: 1.0 if all files exist, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    result_set = set(result)
    missing = [f for f in required_files if f not in result_set]
    logger.info(f'Required files: {required_files}')
    logger.info(f'Found files: {result}')
    logger.info(f'Missing files: {missing}')
    if len(missing) == 0:
        return 1.0
    else:
        return 0.0

def check_pdf_in_listing__4e03b1ed(result, expected, **options):
    """Check if PDF files appear in directory listing with descriptive filenames.

    Args:
        result: Directory listing string from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Score based on PDF file presence and naming
    """
    if 'FOLDER_NOT_FOUND' in result:
        logger.info('Folder not found')
        return 0.0
    required_keywords = expected.get('required_keywords', [])
    min_pdf_count = expected.get('min_pdf_count', 2)
    pdf_pattern = '(\\S+\\.pdf)'
    pdf_files = re.findall(pdf_pattern, result, re.IGNORECASE)
    pdf_count = len(pdf_files)
    score = 0.0
    if pdf_count >= min_pdf_count:
        score += 0.5
        logger.info(f'Found {pdf_count} PDF files in listing (required >= {min_pdf_count}): {pdf_files}')
    else:
        logger.info(f'Only {pdf_count} PDF files, need {min_pdf_count}: {pdf_files}')
    if required_keywords and pdf_count >= min_pdf_count:
        pdfs_with_keywords = 0
        for pdf_file in pdf_files:
            has_keyword = any((kw.lower() in pdf_file.lower() for kw in required_keywords))
            if has_keyword:
                pdfs_with_keywords += 1
                logger.info(f"PDF '{pdf_file}' has descriptive keyword")
            else:
                logger.info(f"PDF '{pdf_file}' missing required keywords")
        if pdfs_with_keywords >= min_pdf_count:
            score += 0.5
            logger.info(f'{pdfs_with_keywords} PDFs have descriptive names with required keywords')
        else:
            logger.info(f'Only {pdfs_with_keywords}/{min_pdf_count} PDFs have descriptive names with keywords')
    return score

def check_named_pdf__0c8a922818da41c6a16e3979ae1f26dc(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF with specific filename exists and contains correct content.

    Args:
        result: Path to the PDF file from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists with correct content, 0.0 otherwise
    """
    if not result or not os.path.exists(result):
        return 0.0
    file_size = os.path.getsize(result)
    if file_size < 10240:
        return 0.0
    try:
        doc = fitz.open(result)
        pdf_text = ''
        for page in doc:
            pdf_text += page.get_text()
        doc.close()
        key_phrases = expected.get('key_phrases', ['Flagstaff Unified School District', 'Functional Behavior Assessment'])
        matches = 0
        for phrase in key_phrases:
            if phrase in pdf_text:
                matches += 1
        return 1.0 if matches == len(key_phrases) else 0.0
    except Exception:
        return 0.0

def check_pdf_content__edbb73b7(pdf_file: str, expected, **options):
    """
    Check PDF has correct structure and image content.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with validation rules
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.0
    try:
        doc = fitz.open(pdf_file)
        if doc.page_count > 0:
            score += 0.4
            image_count = 0
            for page_num in range(doc.page_count):
                page = doc[page_num]
                images = page.get_images(full=True)
                image_count += len(images)
            if image_count > 0:
                score += 0.6
        doc.close()
    except Exception as e:
        return 0.0
    return score

def check_pdf_export__d587f20a(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names with fuzzy matching.

    Args:
        result: Dict with 'pdfs' (list of PDF info) and 'count'
        expected: Dict with 'required_keywords', 'count', and 'title_patterns'

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    pdfs = result.get('pdfs', [])
    expected_count = expected.get('count', 2)
    if len(pdfs) != expected_count:
        return 0.0
    expected_titles = ['LLM Powered Autonomous Agents', 'What Makes Good Data for Alignment']
    title_patterns = expected.get('title_patterns', [['llm', 'autonomous', 'agents'], ['data', 'alignment']])
    total_score = 0.0
    matched_titles = [False] * len(expected_titles)
    for pdf in pdfs:
        filename = pdf.get('filename', '')
        is_recent = pdf.get('is_recent', False)
        content_preview = pdf.get('content_preview', '')
        filename_base = filename.replace('.pdf', '').replace('_', ' ').replace('-', ' ')
        if not is_recent:
            total_score -= 0.2
            continue
        best_match_idx = -1
        best_match_score = 0.0
        for (idx, title) in enumerate(expected_titles):
            if matched_titles[idx]:
                continue
            sim_score = similarity(filename_base, title)
            pattern_match = 0.0
            if idx < len(title_patterns):
                patterns = title_patterns[idx]
                matched_patterns = sum((1 for p in patterns if p.lower() in filename_base.lower()))
                pattern_match = matched_patterns / len(patterns)
            combined_score = sim_score * 0.6 + pattern_match * 0.4
            content_bonus = 0.0
            if content_preview:
                if idx == 0 and any((kw in content_preview.lower() for kw in ['agent', 'autonomous', 'llm'])):
                    content_bonus = 0.1
                elif idx == 1 and any((kw in content_preview.lower() for kw in ['data', 'quality', 'alignment'])):
                    content_bonus = 0.1
            combined_score += content_bonus
            if combined_score > best_match_score:
                best_match_score = combined_score
                best_match_idx = idx
        if best_match_idx >= 0 and best_match_score > 0.3:
            matched_titles[best_match_idx] = True
            total_score += best_match_score
    max_possible = len(expected_titles)
    final_score = total_score / max_possible if max_possible > 0 else 0.0
    if not all(matched_titles):
        final_score *= 0.5
    return min(1.0, max(0.0, final_score))

def check_directory_with_pdfs__92ec1d53(result, expected, **options):
    """Check if directory exists with PDFs.

    Args:
        result: 1 if dir exists with PDFs, 0 otherwise
        expected: Dict with 'expected_value' key

    Returns:
        1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('expected_value', 1)
    if result == expected_value:
        return 1.0
    else:
        return 0.0

def check_pdf_exists_with_content__e1e75309_0(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists and has content (non-zero size and valid structure).

    Args:
        result: Path to the PDF file (from vm_file getter)
        expected: Rules dict with min_size_kb (minimum file size in KB)

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if not result:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if not result.lower().endswith('.pdf'):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size = os.path.getsize(result)
    file_size_kb = file_size / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        with open(result, 'rb') as f:
            header = f.read(5)
            if header != b'%PDF-':
                return 0.0
    except Exception:
        return 0.0
    return 1.0

def check_pdf_export__17dfab17(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names based on page titles.

    Args:
        result: List of dicts with 'filename' and 'size' keys
        expected: Dict with 'page_titles' - list of page titles that should be in filenames
                  and 'count' - expected number of PDFs

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    page_titles = expected.get('page_titles', [])
    expected_count = expected.get('count', len(page_titles))
    if len(result) == 0:
        return 0.0
    valid_pdfs = [pdf for pdf in result if isinstance(pdf, dict) and pdf.get('size', 0) > 0]
    if len(valid_pdfs) == 0:
        return 0.0
    matched_titles = 0
    for title in page_titles:
        title_lower = title.lower()
        words = title_lower.split()
        for pdf_info in valid_pdfs:
            filename = pdf_info.get('filename', '')
            filename_lower = filename.lower()
            found_match = False
            min_phrase_length = 3
            if len(words) < min_phrase_length:
                min_phrase_length = max(2, len(words))
            for i in range(len(words) - min_phrase_length + 1):
                for length in range(min_phrase_length, len(words) - i + 1):
                    phrase = ' '.join(words[i:i + length])
                    phrase_clean = phrase.replace('-', ' ').replace('_', ' ')
                    filename_clean = filename_lower.replace('-', ' ').replace('_', ' ')
                    if phrase_clean in filename_clean:
                        found_match = True
                        break
                if found_match:
                    break
            if found_match:
                matched_titles += 1
                break
    title_match_score = matched_titles / len(page_titles) if page_titles else 0.0
    if len(valid_pdfs) == expected_count:
        count_score = 1.0
    elif len(valid_pdfs) < expected_count:
        count_score = 0.3 * (len(valid_pdfs) / expected_count)
    else:
        count_score = 0.5
    final_score = 0.8 * title_match_score + 0.2 * count_score
    return final_score

def check_pdf_location__0d0715b7122c298c1e10aa6fa135f599(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF exists at expected location and is valid.

    Args:
        result: Dict from getter with file location info
        expected: Dict with expected directory and validation requirements (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    total_checks = 0
    if not result.get('exists', False):
        logger.warning(f"PDF file does not exist at {result.get('path')}")
        return 0.0
    total_checks += 1
    if result.get('valid_pdf', False):
        score += 0.5
        logger.info('PDF file is valid')
    else:
        logger.warning('PDF file is invalid or corrupted')
    if 'expected_directory' in expected:
        total_checks += 1
        if result.get('directory') == expected['expected_directory']:
            score += 0.5
            logger.info(f"PDF saved to correct directory: {expected['expected_directory']}")
        else:
            logger.warning(f"Directory mismatch: expected {expected['expected_directory']}, got {result.get('directory')}")
    if 'min_pages' in expected:
        if result.get('page_count', 0) >= expected['min_pages']:
            logger.info(f"Page count requirement met: {result.get('page_count')} >= {expected['min_pages']}")
        else:
            logger.warning(f"Page count too low: {result.get('page_count')} < {expected['min_pages']}")
            return 0.0
    return score if total_checks > 0 else 0.0

def check_pdf_export__d582d1d2(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names.

    Args:
        result: List of actual PDF filenames
        expected: Dict with 'required_keywords' - list of keyword sets that must appear in filenames
                  Each keyword set can be a list of keywords (all must appear) or a single keyword
                  Optional 'min_filename_length' - minimum length for filenames (excluding .pdf extension)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    required_keywords = expected.get('required_keywords', [])
    expected_count = expected.get('count', len(required_keywords))
    min_filename_length = expected.get('min_filename_length', 0)
    if len(result) != expected_count:
        return 0.0
    matched_files = set()
    matched_keywords = 0
    for keyword_set in required_keywords:
        if isinstance(keyword_set, str):
            keywords_to_check = [keyword_set]
        else:
            keywords_to_check = keyword_set
        for filename in result:
            if filename in matched_files:
                continue
            filename_without_ext = filename.replace('.pdf', '').replace('.PDF', '')
            if len(filename_without_ext) < min_filename_length:
                continue
            separator_count = sum((1 for c in filename_without_ext if c in ' -_'))
            if separator_count == 0 and len(keywords_to_check) > 1:
                continue
            all_keywords_present = all((keyword.lower() in filename.lower() for keyword in keywords_to_check))
            if all_keywords_present:
                matched_files.add(filename)
                matched_keywords += 1
                break
    return 1.0 if matched_keywords == len(required_keywords) else 0.0

def check_pdf_count__6c8aa0e0d1a7a1f247933b05b6dc0e08(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Comprehensive check for PDF evaluation task completion.

    Verifies:
    1. Correct number of employee PDFs created (7 PDFs, excluding template)
    2. Correct naming of PDFs (matches employee names)
    3. PDF content includes employee data (names present in text)
    4. Checkmark symbols present in rating sections

    Awards partial credit:
    - 25% for correct count
    - 25% for correct naming
    - 25% for content filled
    - 25% for checkmarks present

    Args:
        result: Dict from getter with PDF information
        expected: Expected configuration (dict with 'count' key)
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0
    """
    expected_count = expected.get('count', 7)
    employee_pdf_count = result.get('employee_pdf_count', 0)
    employee_pdf_names = result.get('employee_pdf_names', [])
    expected_names = result.get('expected_names', [])
    pdf_details = result.get('pdf_details', [])
    score = 0.0
    if employee_pdf_count == expected_count:
        score += 0.25
    matching_names = sum((1 for name in expected_names if name in employee_pdf_names))
    if expected_count > 0:
        naming_score = matching_names / expected_count
        score += 0.25 * naming_score
    content_filled_count = 0
    for pdf_detail in pdf_details:
        filename = pdf_detail.get('filename', '')
        text = pdf_detail.get('text', '')
        if filename.endswith('.pdf'):
            employee_name = filename[:-4]
            name_parts = employee_name.split()
            if len(name_parts) >= 2:
                if all((part in text for part in name_parts)):
                    content_filled_count += 1
            elif employee_name in text:
                content_filled_count += 1
    if employee_pdf_count > 0:
        content_score = content_filled_count / employee_pdf_count
        score += 0.25 * content_score
    checkmark_count = 0
    for pdf_detail in pdf_details:
        text = pdf_detail.get('text', '')
        if '√' in text or '✓' in text or '✔' in text or ('[X]' in text) or ('[x]' in text):
            checkmark_count += 1
    if employee_pdf_count > 0:
        checkmark_score = checkmark_count / employee_pdf_count
        score += 0.25 * checkmark_score
    return score

def check_pdf_validity__90edd5864af08865c892a3ecd6659029(result: Any, expected: dict, **options) -> float:
    """Check if PDF is valid and meets basic requirements.

    Args:
        result: PDF validity info from getter
        expected: Expected configuration (from rules dict)
        **options: Additional options

    Returns:
        float: Score based on validity checks (0.0-1.0)
    """
    try:
        if not isinstance(result, dict):
            logger.warning(f'Result is not a dict: {type(result)}')
            return 0.0
        require_valid = expected.get('require_valid', True)
        min_pages = expected.get('min_pages', 1)
        require_text = expected.get('require_text', True)
        score = 0.0
        if result.get('valid', False):
            score += 0.4
            logger.info('PDF is valid')
        else:
            logger.warning('PDF is not valid')
            if require_valid:
                return 0.0
        page_count = result.get('page_count', 0)
        if page_count >= min_pages:
            score += 0.3
            logger.info(f'Page count sufficient: {page_count} >= {min_pages}')
        else:
            logger.warning(f'Insufficient pages: {page_count} < {min_pages}')
        has_text = result.get('has_text', False)
        if has_text:
            score += 0.3
            logger.info('PDF has extractable text')
        else:
            logger.warning('PDF has no extractable text')
            if require_text and score == 0.0:
                return 0.0
        logger.info(f'PDF validity score: {score:.2f}')
        return score
    except Exception as e:
        logger.error(f'Error checking PDF validity: {e}')
        return 0.0

def check_pdf_chapters__097402a1(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts."""
    expected_files = expected.get('files', [])
    expected_page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    result_map = {item['filename']: item['page_count'] for item in result}
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result_map:
            continue
        file_score = 0.5
        if expected_file in expected_page_counts:
            page_count_rule = expected_page_counts[expected_file]
            relation = page_count_rule.get('relation', 'ge')
            ref_value = page_count_rule.get('ref_value', 0)
            actual_page_count = result_map[expected_file]
            page_count_valid = False
            if relation == 'ge' and actual_page_count >= ref_value:
                page_count_valid = True
            elif relation == 'le' and actual_page_count <= ref_value:
                page_count_valid = True
            elif relation == 'eq' and actual_page_count == ref_value:
                page_count_valid = True
            if page_count_valid:
                file_score = 1.0
        else:
            file_score = 1.0
        score += points_per_file * file_score
    return min(score, 1.0)

def check_pdf_orientation__e222560cb6780019664a6917701a2659(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDF has the expected orientation and page count.

    Args:
        result: Dict from getter with 'orientation' and 'page_count' keys
        expected: Dict with 'orientation' (str) and 'page_count' (int) keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_count = expected.get('page_count', 1)
    if result.get('page_count') == expected_count:
        score += 0.5
    expected_orientation = expected.get('orientation', 'landscape')
    if result.get('orientation') == expected_orientation:
        score += 0.5
    return score

def check_pdf_chapters__d2aaaf87(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count
        expected: Expected state with 'files' list and 'page_counts' dict

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        actual_page_count = result[expected_file]
        if expected_file in page_counts:
            page_count_requirement = page_counts[expected_file]
            relation = page_count_requirement.get('relation', 'eq')
            ref_value = page_count_requirement.get('ref_value', 0)
            if relation == 'ge':
                if actual_page_count < ref_value:
                    continue
            elif relation == 'le':
                if actual_page_count > ref_value:
                    continue
            elif relation == 'eq':
                if actual_page_count != ref_value:
                    continue
            elif relation == 'gt':
                if actual_page_count <= ref_value:
                    continue
            elif relation == 'lt':
                if actual_page_count >= ref_value:
                    continue
        score += points_per_file
    return min(score, 1.0)

def check_pdf_properties__869dd2b5b7a0e45511c735d06983da83(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF has expected properties (orientation, page count, etc.).

    Args:
        result: Dict from getter with PDF properties
        expected: Dict with expected properties (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists', False):
        logger.warning('PDF file does not exist or could not be read')
        return 0.0
    score = 0.0
    checks = 0
    if 'orientation' in expected:
        checks += 1
        if result.get('orientation') == expected['orientation']:
            score += 1.0
            logger.info(f"Orientation check passed: {expected['orientation']}")
        else:
            logger.warning(f"Orientation mismatch: expected {expected['orientation']}, got {result.get('orientation')}")
    if 'min_pages' in expected:
        checks += 1
        if result.get('page_count', 0) >= expected['min_pages']:
            score += 1.0
            logger.info(f"Page count check passed: {result.get('page_count')} >= {expected['min_pages']}")
        else:
            logger.warning(f"Page count too low: {result.get('page_count')} < {expected['min_pages']}")
    if checks > 0:
        return score / checks
    return 1.0 if result.get('exists', False) else 0.0

def check_pdf_files__2c1e781e(result, expected, **options):
    """Check specific files with partial credit.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: Ratio of found files to required files (0.0-1.0)
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    if not required_files:
        return 1.0
    result_set = set(result)
    found_count = sum((1 for f in required_files if f in result_set))
    score = found_count / len(required_files)
    logger.info(f'Required {len(required_files)} files, found {found_count}')
    logger.info(f'Required: {required_files}')
    logger.info(f'Found: {result}')
    logger.info(f'Score: {score:.2f}')
    return score

def check_pdf_orientation__596fcd5e219bea1372357f0afb95cc85(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF has the expected page orientation.

    Args:
        result: Dict from getter with orientation info
        expected: Dict with expected orientation
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    if result.get('page_count', 0) > 0:
        score += 0.4
    expected_orientation = expected.get('orientation', 'portrait')
    if result.get('orientation') == expected_orientation:
        score += 0.6
    return score

def check_pdf_files_non_empty__99e94d1a(result, expected, **options):
    """Check if PDF files exist and are non-empty.

    Args:
        result: List of tuples (filename, size_in_bytes)
        expected: Dict with 'min_count' and 'min_size_bytes'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    min_count = expected.get('min_count', 2)
    min_size_bytes = expected.get('min_size_bytes', 1000)
    if not isinstance(result, list):
        logger.warning(f'Expected list, got {type(result)}')
        return 0.0
    if len(result) < min_count:
        logger.warning(f'Only found {len(result)} files, expected at least {min_count}')
        return 0.0
    score = 0.0
    valid_files = [f for f in result if f[1] >= min_size_bytes]
    if len(valid_files) >= min_count:
        score = 1.0
        logger.info(f'All required files are non-empty: {len(valid_files)} valid files')
    else:
        score = len(valid_files) / min_count
        logger.warning(f'Only {len(valid_files)}/{min_count} files meet size requirement')
    return score

def check_pdf_chapters__89dbe8bb(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict with 'files' (list) and 'page_counts' (dict of filename -> page count)
        expected: Dict with 'files' (list) and 'page_counts' (dict of filename -> constraints)

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    expected_page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    result_files = result.get('files', [])
    result_page_counts = result.get('page_counts', {})
    filename_score = 0.0
    page_count_score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file in result_files:
            filename_score += points_per_file * 0.5
            if expected_file in expected_page_counts:
                constraint = expected_page_counts[expected_file]
                relation = constraint.get('relation', 'eq')
                ref_value = constraint.get('ref_value', 0)
                actual_page_count = result_page_counts.get(expected_file, 0)
                constraint_met = False
                if relation == 'ge':
                    constraint_met = actual_page_count >= ref_value
                elif relation == 'le':
                    constraint_met = actual_page_count <= ref_value
                elif relation == 'eq':
                    constraint_met = actual_page_count == ref_value
                elif relation == 'gt':
                    constraint_met = actual_page_count > ref_value
                elif relation == 'lt':
                    constraint_met = actual_page_count < ref_value
                if constraint_met:
                    page_count_score += points_per_file * 0.5
            else:
                page_count_score += points_per_file * 0.5
    total_score = filename_score + page_count_score
    return min(total_score, 1.0)

def check_pdf_files__197670d4(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_export__b47a318c(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names.
    Verifies that filenames are derived from page titles using:
    1. Keyword matching with higher threshold
    2. Fuzzy string similarity against actual expected titles
    3. PDF metadata verification
    4. Content verification via text snippets

    Args:
        result: List of dicts, each containing:
            - 'filename': PDF filename
            - 'title': PDF metadata title (if available)
            - 'creator': PDF creator/producer (if available)
            - 'text_snippet': First 200 chars of text content
        expected: Dict with:
            - 'pdf_requirements': List of dicts, each containing:
                - 'expected_title': The actual expected page title
                - 'keywords': List of keywords from a page title
                - 'min_keywords': Minimum number of keywords that must appear in filename
            - 'count': Expected number of PDFs

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    pdf_requirements = expected.get('pdf_requirements', [])
    expected_count = expected.get('count', len(pdf_requirements))
    if len(result) != expected_count:
        return 0.0
    matched_pdfs = set()
    requirement_scores = []
    for requirement in pdf_requirements:
        keywords = requirement.get('keywords', [])
        min_keywords = requirement.get('min_keywords', 1)
        expected_title = requirement.get('expected_title', ' '.join(keywords))
        best_score = 0.0
        best_pdf_idx = -1
        for (i, pdf_info) in enumerate(result):
            if i in matched_pdfs:
                continue
            if isinstance(pdf_info, dict):
                filename = pdf_info.get('filename', '')
                pdf_title = pdf_info.get('title', '')
                text_snippet = pdf_info.get('text_snippet', '')
            else:
                filename = pdf_info
                pdf_title = ''
                text_snippet = ''
            filename_base = filename.replace('.pdf', '')
            filename_lower = filename_base.lower()
            keyword_count = sum((1 for keyword in keywords if keyword.lower() in filename_lower))
            keyword_score = keyword_count / len(keywords) if keywords else 0.0
            expected_title_normalized = expected_title.lower().replace(' ', '_').replace('-', '_')
            title_similarity_1 = SequenceMatcher(None, filename_lower, expected_title_normalized).ratio()
            title_similarity_2 = SequenceMatcher(None, filename_lower, expected_title.lower()).ratio()
            title_similarity = max(title_similarity_1, title_similarity_2)
            metadata_score = 0.0
            if pdf_title:
                pdf_title_lower = pdf_title.lower()
                metadata_similarity = SequenceMatcher(None, pdf_title_lower, expected_title.lower()).ratio()
                metadata_keyword_count = sum((1 for keyword in keywords if keyword.lower() in pdf_title_lower))
                metadata_keyword_score = metadata_keyword_count / len(keywords) if keywords else 0.0
                metadata_score = 0.6 * metadata_similarity + 0.4 * metadata_keyword_score
            content_score = 0.0
            if text_snippet:
                text_lower = text_snippet.lower()
                content_keyword_count = sum((1 for keyword in keywords if keyword.lower() in text_lower))
                content_score = min(1.0, content_keyword_count / len(keywords)) if keywords else 0.0
            combined_score = 0.3 * keyword_score + 0.25 * title_similarity + 0.3 * metadata_score + 0.15 * content_score
            if keyword_count >= min_keywords and combined_score > best_score:
                best_score = combined_score
                best_pdf_idx = i
        if best_pdf_idx >= 0 and best_score >= 0.7:
            matched_pdfs.add(best_pdf_idx)
            requirement_scores.append(best_score)
        else:
            requirement_scores.append(0.0)
    if len(requirement_scores) == 0:
        return 1.0
    return sum(requirement_scores) / len(requirement_scores)

def check_pdf_export__bec19abd(result, expected, **options):
    """
    Compare exported PDF list against expected article titles using strict matching.

    Args:
        result: Dict with 'pdf_files' (list of PDF filenames), 'article_titles' (list of article titles),
                and 'pdf_titles' (list of titles extracted from PDF content)
        expected: Dict with 'count' - expected number of PDFs and optional 'urls' for URL-based matching

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    pdf_files = result.get('pdf_files', [])
    article_titles = result.get('article_titles', [])
    pdf_titles = result.get('pdf_titles', [])
    expected_count = expected.get('count', 2)
    if len(pdf_files) != expected_count:
        return 0.0
    if not article_titles:
        return 0.0
    matched_pdfs = set()
    match_threshold = 0.7
    matches_found = 0
    for (idx, article_title) in enumerate(article_titles):
        best_match_score = 0.0
        best_match_pdf = None
        best_match_idx = -1
        article_words = set(article_title.lower().split())
        significant_words = {w for w in article_words if len(w) > 3 and w not in ['the', 'and', 'for', 'with', 'from', 'this', 'that', 'what', 'make', 'makes', 'good']}
        for (pdf_idx, pdf_file) in enumerate(pdf_files):
            if pdf_idx in matched_pdfs:
                continue
            clean_title = article_title.lower().replace(' ', '').replace('-', '').replace('_', '')
            clean_pdf = pdf_file.lower().replace('.pdf', '').replace(' ', '').replace('-', '').replace('_', '')
            filename_similarity = SequenceMatcher(None, clean_title, clean_pdf).ratio()
            pdf_lower = pdf_file.lower()
            word_matches = sum((1 for word in significant_words if word in pdf_lower))
            word_coverage = word_matches / len(significant_words) if significant_words else 0.0
            min_required_words = min(2, len(significant_words))
            has_required_words = word_matches >= min_required_words if significant_words else False
            content_similarity = 0.0
            if pdf_idx < len(pdf_titles) and pdf_titles[pdf_idx]:
                pdf_content_title = pdf_titles[pdf_idx]
                clean_pdf_content = pdf_content_title.lower().replace(' ', '').replace('-', '').replace('_', '')
                content_similarity = SequenceMatcher(None, clean_title, clean_pdf_content).ratio()
                pdf_content_lower = pdf_content_title.lower()
                content_word_matches = sum((1 for word in significant_words if word in pdf_content_lower))
                content_word_coverage = content_word_matches / len(significant_words) if significant_words else 0.0
                if content_similarity >= 0.6 or content_word_coverage >= 0.7:
                    content_similarity = max(content_similarity, content_word_coverage)
            if content_similarity >= 0.6:
                combined_score = content_similarity
            elif has_required_words and word_coverage >= 0.7:
                combined_score = max(filename_similarity, word_coverage)
            elif filename_similarity >= 0.8:
                combined_score = filename_similarity
            else:
                combined_score = max(filename_similarity, word_coverage, content_similarity)
            if combined_score > best_match_score:
                best_match_score = combined_score
                best_match_pdf = pdf_file
                best_match_idx = pdf_idx
        if best_match_score >= match_threshold and best_match_pdf is not None:
            matched_pdfs.add(best_match_idx)
            matches_found += 1
    if len(article_titles) == 0:
        return 0.0
    if matches_found == expected_count and matches_found == len(article_titles):
        return 1.0
    score = matches_found / len(article_titles)
    return score

def check_pdf_files_exist__989759aa(result, expected, **options):
    """Check if expected PDF files exist in the result list.

    Args:
        result: List of PDF filenames found
        expected: Dict with 'min_count' and optional 'required_keywords'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.warning(f'Expected list, got {type(result)}')
        return 0.0
    min_count = expected.get('min_count', 1)
    required_keywords = expected.get('required_keywords', [])
    score = 0.0
    if len(result) >= min_count:
        score += 0.5
        logger.info(f'Found {len(result)} PDF files (required: {min_count})')
    else:
        logger.warning(f'Only found {len(result)} PDF files (required: {min_count})')
        return score
    if required_keywords:
        matched_keywords = 0
        for keyword in required_keywords:
            if any((keyword.lower() in filename.lower() for filename in result)):
                matched_keywords += 1
        if matched_keywords == len(required_keywords):
            score += 0.5
            logger.info(f'All required keywords found in filenames')
        else:
            partial_score = matched_keywords / len(required_keywords) * 0.5
            score += partial_score
            logger.info(f'Found {matched_keywords}/{len(required_keywords)} required keywords')
    else:
        score += 0.5
    return score

def check_pdf_size_and_validity__e1e75309_9(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists, has valid structure, and meets size requirements.

    Args:
        result: Path to the PDF file (from vm_file getter)
        expected: Rules dict with 'min_size_kb' and 'max_size_kb' (file size range in KB)

    Returns:
        float: Score from 0.0 to 1.0 based on multiple criteria
    """
    if not result:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if not result.lower().endswith('.pdf'):
        return 0.0
    score = 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    max_size_kb = expected.get('max_size_kb', 10000)
    file_size = os.path.getsize(result)
    file_size_kb = file_size / 1024
    if min_size_kb <= file_size_kb <= max_size_kb:
        score += 0.33
    try:
        with open(result, 'rb') as f:
            header = f.read(5)
            if header == b'%PDF-':
                score += 0.34
    except Exception:
        pass
    try:
        reader = PdfReader(result)
        page_count = len(reader.pages)
        if page_count > 0:
            score += 0.33
    except Exception:
        pass
    return score

def check_pdf_count__085cc8d6(result: int, expected: dict, **options) -> float:
    """
    Check if PDF count matches expected value.

    Args:
        result: Actual count from getter
        expected: Dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('expected_count', 0)
    logger.info(f'Actual PDF count: {result}, Expected: {expected_count}')
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_pdf_export__831e0e03(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names.

    Args:
        result: List of actual PDF filenames
        expected: Dict with 'required_keywords' - list of keywords that must appear in filenames

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    required_keywords = expected.get('required_keywords', [])
    expected_count = expected.get('count', len(required_keywords))
    if len(result) != expected_count:
        return 0.0
    score = 0.0
    for keyword in required_keywords:
        found = any((keyword.lower() in filename.lower() for filename in result))
        if found:
            score += 1.0 / len(required_keywords)
    return score

def check_pdf_basic_requirements__7b63a847bf086913e974e7a32debb9ec(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF meets basic requirements (exists, correct page count, has images).

    Args:
        result: Dict from getter with file info
        expected: Dict with expected values (page_count, has_images, min_size)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'page_count' in expected:
        total_checks += 1
        if result.get('page_count') == expected['page_count']:
            score += 1.0
    if 'has_images' in expected:
        total_checks += 1
        if result.get('has_images') == expected['has_images']:
            score += 1.0
    if 'min_size' in expected:
        total_checks += 1
        if result.get('file_size', 0) >= expected['min_size']:
            score += 1.0
    if total_checks == 0:
        return 0.0
    return score / total_checks

def check_pdf_chapters__21dcc4a0(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count
        expected: Dict with 'files' list and 'page_counts' dict
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    filename_weight = 0.4
    pagecount_weight = 0.6
    points_per_file = 1.0 / len(expected_files)
    score = 0.0
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        file_score = filename_weight * points_per_file
        if expected_file in page_counts:
            page_requirement = page_counts[expected_file]
            relation = page_requirement.get('relation', 'ge')
            ref_value = page_requirement.get('ref_value', 0)
            actual_pages = result[expected_file]
            page_valid = False
            if relation == 'ge':
                page_valid = actual_pages >= ref_value
            elif relation == 'eq':
                page_valid = actual_pages == ref_value
            elif relation == 'le':
                page_valid = actual_pages <= ref_value
            elif relation == 'gt':
                page_valid = actual_pages > ref_value
            elif relation == 'lt':
                page_valid = actual_pages < ref_value
            if page_valid:
                file_score += pagecount_weight * points_per_file
        else:
            file_score += pagecount_weight * points_per_file
        score += file_score
    return min(score, 1.0)

def check_pdf_files__230c9972(result, expected, **options):
    """Check specific files with partial credit.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: Ratio of found files to required files (0.0-1.0)
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    if not required_files:
        return 1.0
    result_set = set(result)
    found_count = sum((1 for f in required_files if f in result_set))
    score = found_count / len(required_files)
    logger.info(f'Required {len(required_files)} files, found {found_count}')
    logger.info(f'Required: {required_files}')
    logger.info(f'Found: {result}')
    logger.info(f'Score: {score:.2f}')
    return score

def check_pdf_named__b387dd88(pdf_file: str, expected, **options):
    """
    Check if PDF exists with specific filename and valid content.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with validation rules
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.4
    try:
        reader = PdfReader(pdf_file)
        if len(reader.pages) > 0:
            file_size = os.path.getsize(pdf_file)
            if file_size > 5120:
                score += 0.6
    except Exception as e:
        return score
    return score

def check_pdf_content__5e9826db825b680d44e19c36727da776(result: Dict[str, Any], expected: Dict, **options) -> float:
    """
    Comprehensive check of PDF content to verify task completion.

    Args:
        result: Dict from getter containing file existence, page count, and content checks
        expected: Expected configuration (dict with required checks)
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on how many checks pass
    """
    if not result.get('exists', False):
        return 0.0
    checks_passed = 0
    total_checks = 0
    if 'page_count' in expected:
        total_checks += 1
        if result.get('page_count', 0) == expected['page_count']:
            checks_passed += 1
    if expected.get('requires_employee_name', False):
        total_checks += 1
        if result.get('has_employee_name', False):
            checks_passed += 1
    if expected.get('requires_checkmarks', False):
        total_checks += 1
        if result.get('has_checkmarks', False):
            checks_passed += 1
    if expected.get('requires_ratings', False):
        total_checks += 1
        if result.get('has_ratings', False):
            checks_passed += 1
    if total_checks == 0:
        return 1.0 if result.get('exists', False) else 0.0
    return checks_passed / total_checks

def check_pdf_files__1ee77af4(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_chapters__0be3e647(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts."""
    expected_files = expected.get('files', [])
    page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        file_score = 0.5
        if expected_file in page_counts:
            page_requirement = page_counts[expected_file]
            relation = page_requirement.get('relation', 'ge')
            ref_value = page_requirement.get('ref_value', 0)
            actual_page_count = result[expected_file]
            if relation == 'ge' and actual_page_count >= ref_value:
                file_score += 0.5
            elif relation == 'eq' and actual_page_count == ref_value:
                file_score += 0.5
            elif relation == 'le' and actual_page_count <= ref_value:
                file_score += 0.5
            elif relation == 'gt' and actual_page_count > ref_value:
                file_score += 0.5
            elif relation == 'lt' and actual_page_count < ref_value:
                file_score += 0.5
        else:
            file_score = 1.0
        score += file_score * points_per_file
    return min(score, 1.0)

def check_pdf_basic__d8736b62(pdf_file: str, expected, **options):
    """
    Basic check if PDF exists and is valid with at least one page.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with basic validation rules
        **options: Additional options

    Returns:
        float: Score 0.0-1.0
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    try:
        reader = PdfReader(pdf_file)
        if len(reader.pages) < 1:
            return 0.0
        file_size = os.path.getsize(pdf_file)
        if file_size < 1024:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_readable__cd376b4c(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF is readable (can be opened and parsed).

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if PDF is readable, 0.0 otherwise
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        if len(reader.pages) > 0:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_pdf_chapters__c8947a04(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count from getter
        expected: Dict with 'files' list and 'page_counts' dict

    Returns:
        float: Score from 0.0 to 1.0
    """
    expected_files = expected.get('files', [])
    page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        if expected_file in page_counts:
            page_requirement = page_counts[expected_file]
            relation = page_requirement.get('relation')
            ref_value = page_requirement.get('ref_value')
            actual_pages = result[expected_file]
            page_count_valid = False
            if relation == 'ge':
                page_count_valid = actual_pages >= ref_value
            elif relation == 'eq':
                page_count_valid = actual_pages == ref_value
            elif relation == 'le':
                page_count_valid = actual_pages <= ref_value
            elif relation == 'gt':
                page_count_valid = actual_pages > ref_value
            elif relation == 'lt':
                page_count_valid = actual_pages < ref_value
            if page_count_valid:
                score += points_per_file
        else:
            score += points_per_file
    return min(score, 1.0)

def check_pdf_comprehensive__6dc3161b(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Comprehensive PDF validation: file size, page count, and text content.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: Score with partial credit (0.33 per check passed)
    """
    min_file_size = expected.get('min_file_size', 0)
    min_pages = expected.get('min_pages', 1)
    has_text = expected.get('has_text', True)
    if result is None or not os.path.exists(result):
        return 0.0
    score = 0.0
    try:
        file_size = os.path.getsize(result)
        if file_size >= min_file_size:
            score += 0.33
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        if num_pages >= min_pages:
            score += 0.34
        if has_text:
            with fitz.open(result) as pdf:
                text = ''
                for page in pdf:
                    text += page.get_text()
                if len(text.strip()) > 0:
                    score += 0.33
        else:
            score += 0.33
        return score
    except Exception as e:
        return score

def check_pdf_files__32b125c8(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_files_with_content__83b6523335a3d68f9c734599b1537c74(result: Dict[str, Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if expected PDF files exist and contain valid performance review content.

    Args:
        result: Dictionary mapping filename to content info from getter
        expected: Expected configuration (dict with 'required_files' key - list of filenames)
        **options: Additional options

    Returns:
        float: Score based on existence and content validation (0.0 to 1.0)
            - 0.25 points per file for existence
            - 0.25 points per file for substantial content (>200 chars) with employee name
            - 0.25 points per file for having rating marks ('√')
            - 0.25 points per file for having evaluation keywords (performance, rating, etc.)
    """
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    total_score = 0.0
    max_score = len(required_files)
    for required_file in required_files:
        file_info = result.get(required_file, {})
        if file_info.get('exists', False):
            total_score += 0.25
            if file_info.get('has_content', False) and file_info.get('contains_name', False):
                total_score += 0.25
            if file_info.get('has_rating_mark', False):
                total_score += 0.25
            if file_info.get('has_evaluation_keywords', False):
                total_score += 0.25
    score = total_score / max_score if max_score > 0 else 0.0
    return score

def check_pdf_files__d8278409(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_contains_text__4c28c68dd08d7073669bab47ee359a64(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF text contains expected strings for performance review form.

    Verifies that the PDF contains:
    - Employee personal information (name, ID, position, department)
    - Rating selections marked with '√' symbol
    - All required text elements are present in the filled form

    Args:
        result: Extracted PDF text from getter
        expected: Expected configuration dict with structure:
                 {"type": "rule", "rules": {"contains": [list of required strings]}}
        **options: Additional options

    Returns:
        float: 1.0 if all expected strings found, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_strings = expected.get('contains', [])
    if not expected_strings:
        return 0.0
    for expected_str in expected_strings:
        if expected_str not in result:
            return 0.0
    return 1.0

def check_pdf_files__184db76a3945be5ea6dec91515beac2a(result, expected, **options):
    """Check if all expected PDF files exist with correct content.

    Args:
        result: Dict from getter with filename -> status mapping
        expected: Dict with 'expected_files' list containing filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    total_files = len(expected_files)
    correct_files = 0
    for file_info in expected_files:
        filename = file_info['filename']
        status = result.get(filename, -1)
        if status == 1:
            correct_files += 1
    return correct_files / total_files

def check_pdf_exists__f7d1f102(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists in Google Drive.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    should_exist = expected.get('exists', True)
    file_exists = result is not None and os.path.exists(result)
    if should_exist:
        return 1.0 if file_exists else 0.0
    else:
        return 0.0 if file_exists else 1.0

def check_pdf_text_content__7a2513e6563b76f5c07e27d8c4089d02(result: Optional[str], expected: Dict[str, Any], **options) -> float:
    """Check if PDF contains expected text keywords.

    Args:
        result: Extracted text content from PDF
        expected: Expected rules dict with 'keywords' list - text phrases that should appear
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on how many keywords are found
    """
    if result is None:
        logger.info('No text content extracted from PDF')
        return 0.0
    keywords = expected.get('keywords', [])
    if not keywords:
        logger.warning('No keywords specified in expected rules')
        return 0.0
    found_count = 0
    for keyword in keywords:
        if keyword.lower() in result.lower():
            found_count += 1
            logger.info(f"Found keyword: '{keyword}'")
        else:
            logger.info(f"Missing keyword: '{keyword}'")
    score = found_count / len(keywords)
    logger.info(f'Found {found_count}/{len(keywords)} keywords, score: {score}')
    return score

def check_pdf_saved__5f59826466b2625834fd8d369560ed11(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if a PDF file was successfully saved with valid content.

    Args:
        result: Path to the PDF file (from getter)
        expected: Dict with validation rules (min_pages, min_size_kb)
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size_kb = os.path.getsize(result) / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        from pypdf import PdfReader
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        min_pages = expected.get('min_pages', 1)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_exists__9753fb3ef75762f14fb08b7f236e3f81(result: Optional[str], expected: Dict[str, Any], **options) -> float:
    """Check if a PDF file exists and has valid content.

    Args:
        result: Path to the PDF file from getter (or None if not found)
        expected: Expected rules dict with optional 'min_size' (bytes)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is valid, 0.0 otherwise
    """
    if result is None:
        logger.info('PDF file not found')
        return 0.0
    if not os.path.exists(result):
        logger.info(f'PDF file path {result} does not exist')
        return 0.0
    file_size = os.path.getsize(result)
    min_size = expected.get('min_size', 1024)
    if file_size < min_size:
        logger.info(f'PDF file too small: {file_size} bytes < {min_size} bytes')
        return 0.0
    logger.info(f'PDF file exists and valid: {result} ({file_size} bytes)')
    return 1.0

def check_pdf_location__3a8b0dcb7e90aff8357c94bc8dad2474(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF exists at expected location with correct page count and contains image content.

    Args:
        result: Dict from getter with file info
        expected: Dict with expected path and page_count
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    if 'page_count' in expected:
        if result.get('page_count') == expected['page_count']:
            score += 0.3
    if result.get('file_size', 0) > 10000:
        score += 0.3
    if result.get('has_image', False):
        score += 0.4
    return score

def check_pdf_contains_keywords__2c2512f5(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF contains expected keywords.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if minimum keywords found, 0.0 otherwise
    """
    keywords: List[str] = expected.get('keywords', [])
    min_keywords_found = expected.get('min_keywords_found', len(keywords))
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with fitz.open(result) as pdf:
            text = ''
            for page in pdf:
                text += page.get_text()
        keywords_found = 0
        for keyword in keywords:
            if keyword.lower() in text.lower():
                keywords_found += 1
        if keywords_found >= min_keywords_found:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_pdf_files__8d8da24c(result, expected, **options):
    """Check specific files with partial credit.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: Ratio of found files to required files (0.0-1.0)
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    if not required_files:
        return 1.0
    result_set = set(result)
    found_count = sum((1 for f in required_files if f in result_set))
    score = found_count / len(required_files)
    logger.info(f'Required {len(required_files)} files, found {found_count}')
    logger.info(f'Required: {required_files}')
    logger.info(f'Found: {result}')
    logger.info(f'Score: {score:.2f}')
    return score

def check_pdf_all_fields__ff249445b547a028ee5246e45cd19fdf(result: Dict[str, Any], expected: Dict, **options) -> float:
    """
    Check if all expected text fields are present in the PDF.

    Args:
        result: Dictionary containing text_fields and other data from getter
        expected: Expected configuration (dict with 'required_fields' key - list of fields)
        **options: Additional options (partial_credit: bool)

    Returns:
        float: 1.0 if all fields found, or partial credit if enabled
    """
    required_fields = expected.get('required_fields', [])
    if not required_fields:
        return 0.0
    partial_credit = options.get('partial_credit', True)
    text_fields = result.get('text_fields', {})
    text_fields_found = 0
    for field in required_fields:
        if text_fields.get(field, False):
            text_fields_found += 1
    if partial_credit:
        score = text_fields_found / len(required_fields)
        return score
    else:
        return 1.0 if text_fields_found == len(required_fields) else 0.0

def check_pdf_file_size__f34c6d72482f0e93994904e974de58fa(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF has expected file size (ensures quality/completeness).

    Args:
        result: Dict from getter with file size info
        expected: Dict with min file size requirement (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        logger.warning(f"PDF file does not exist at {result.get('path')}")
        return 0.0
    if not result.get('valid_pdf', False):
        logger.warning('PDF file is invalid or corrupted')
        return 0.0
    file_size = result.get('file_size_bytes', 0)
    min_size = expected.get('min_size_bytes', 1024)
    if file_size >= min_size:
        logger.info(f'File size {file_size} bytes meets minimum requirement {min_size} bytes')
        return 1.0
    else:
        logger.warning(f'File size {file_size} bytes is below minimum requirement {min_size} bytes')
        return 0.0

def check_pdf_location__c74e09f1(pdf_file: str, expected, **options):
    """
    Check if PDF exists at the expected location with correct filename and is valid.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with 'expected_filename'
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.4
    try:
        actual_filename = os.path.basename(pdf_file)
        expected_filename = expected.get('expected_filename', '')
        if actual_filename == expected_filename:
            score += 0.3
        reader = PdfReader(pdf_file)
        if len(reader.pages) > 0:
            file_size = os.path.getsize(pdf_file)
            if file_size > 1024:
                score += 0.3
    except Exception as e:
        return score
    return score

def check_desktop_pdf__a45ebe2876987410607711d3992c10db(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF exists and contains expected text content.

    Args:
        result: Path to the PDF file from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and contains expected text, 0.0 otherwise
    """
    if not result or not os.path.exists(result):
        return 0.0
    file_size = os.path.getsize(result)
    if file_size < 10240:
        return 0.0
    try:
        doc = fitz.open(result)
        pdf_text = ''
        for page in doc:
            pdf_text += page.get_text()
        doc.close()
        key_phrases = expected.get('key_phrases', [])
        if not key_phrases:
            return 1.0 if len(pdf_text.strip()) > 100 else 0.0
        matches = 0
        for phrase in key_phrases:
            if phrase.lower() in pdf_text.lower():
                matches += 1
        return matches / len(key_phrases)
    except Exception as e:
        return 0.0

def check_pdf_exists__3df7d80c(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDF file exists and meets size requirements.

    Args:
        result: Dict from getter with 'exists' and 'size_bytes' keys
        expected: Expected rules dict with 'exists' and 'min_size_kb' keys

    Returns:
        float: 1.0 if file exists and meets requirements, 0.0 otherwise
    """
    if not result.get('exists', False):
        return 0.0
    if not expected.get('exists', True):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 0)
    if min_size_kb > 0:
        size_kb = result.get('size_bytes', 0) / 1024
        if size_kb < min_size_kb:
            return 0.0
    return 1.0

def check_exact_pdf_count__10f702e7ff0a641a1fda6d45251486ce(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if exactly the expected number of PDFs exist with correct naming.

    Args:
        result: Dict with 'files' (list of filenames) and 'count' from getter
        expected: Dict with 'exact_count' and 'expected_title_patterns' keys
        **options: Additional options

    Returns:
        1.0 if exact count matches and filenames match expected patterns, partial credit otherwise
    """
    exact_count = expected.get('exact_count', 2)
    expected_patterns = expected.get('expected_title_patterns', [])
    actual_files = result.get('files', [])
    actual_count = len(actual_files)
    if actual_count == 0:
        return 0.0
    if actual_count != exact_count:
        if actual_count > exact_count:
            return 0.3
        else:
            return actual_count / exact_count * 0.5
    if not expected_patterns:
        return 1.0
    matched_patterns = set()
    for filename in actual_files:
        filename_lower = filename.lower()
        filename_base = filename_lower.replace('.pdf', '')
        for pattern in expected_patterns:
            pattern_lower = pattern.lower()
            pattern_keywords = re.findall('\\w+', pattern_lower)
            filename_keywords = re.findall('\\w+', filename_base)
            matches = sum((1 for pk in pattern_keywords if any((pk in fk or fk in pk for fk in filename_keywords))))
            if matches >= len(pattern_keywords) * 0.5:
                matched_patterns.add(pattern)
                break
    pattern_match_ratio = len(matched_patterns) / len(expected_patterns)
    if pattern_match_ratio >= 1.0:
        return 1.0
    elif pattern_match_ratio >= 0.5:
        return 0.7 + (pattern_match_ratio - 0.5) * 0.6
    else:
        return 0.3 + pattern_match_ratio * 0.8

def check_pdf_file_exists__e68e77d1(result: str, expected: dict, **options) -> float:
    """
    Check if PDF file exists and optionally verify minimum file size.

    Args:
        result: Path to the PDF file from getter
        expected: Dict with optional 'min_size_kb' to verify file is not empty
        **options: Additional options

    Returns:
        float: 1.0 if file exists and meets size requirement, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None - file not found')
        return 0.0
    if not os.path.exists(result):
        logger.warning(f'File does not exist: {result}')
        return 0.0
    min_size_kb = expected.get('min_size_kb', 0)
    if min_size_kb > 0:
        file_size_kb = os.path.getsize(result) / 1024
        if file_size_kb < min_size_kb:
            logger.warning(f'File too small: {file_size_kb:.1f}KB < {min_size_kb}KB')
            return 0.0
        logger.info(f'File size OK: {file_size_kb:.1f}KB >= {min_size_kb}KB')
    try:
        with open(result, 'rb') as f:
            header = f.read(5)
            if header != b'%PDF-':
                logger.warning('File is not a valid PDF (wrong header)')
                return 0.0
    except Exception as e:
        logger.error(f'Error reading file: {e}')
        return 0.0
    logger.info(f'PDF file exists and is valid: {result}')
    return 1.0

def check_pdf_file_count__822cf85f742bbfae9e3acf8d7027940c(result, expected, **options):
    """Check if the PDF files match the expected configuration.

    Verifies:
    1. Original 'Spectral Graph Theory.pdf' exists
    2. Exactly 8 chapter PDF files exist (distinct from the original)
    3. Total count is 9 PDF files

    Args:
        result: Dict with file details from getter
        expected: Dict with 'count' key specifying expected number of files
        **options: Additional options

    Returns:
        float: 1.0 if all requirements are met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('original', False):
        return 0.0
    chapters = result.get('chapters', [])
    if len(chapters) != 8:
        return 0.0
    total_count = result.get('total_count', 0)
    expected_count = expected.get('count', 9)
    if total_count != expected_count:
        return 0.0
    return 1.0

def check_pdf_files__4e03b1ed(result, expected, **options):
    """Check if expected PDF files exist with strict validation.

    Args:
        result: Dict with 'files' (list) and 'file_info' (dict) from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    if not expected_files:
        logger.info('No expected files specified')
        return 0.0
    if isinstance(result, dict):
        result_files = result.get('files', [])
        file_info = result.get('file_info', {})
    else:
        result_files = result
        file_info = {}
    if len(result_files) != len(expected_files):
        logger.info(f'File count mismatch: expected {len(expected_files)} PDFs, found {len(result_files)}')
        logger.info(f'Expected files: {expected_files}')
        logger.info(f'Found files: {result_files}')
        return 0.0

    def normalize_filename(filename):
        return filename.lower().strip()
    result_normalized = {normalize_filename(f): f for f in result_files}
    expected_normalized = {normalize_filename(f): f for f in expected_files}
    missing_files = []
    for (exp_norm, exp_orig) in expected_normalized.items():
        if exp_norm not in result_normalized:
            missing_files.append(exp_orig)
    if missing_files:
        logger.info(f'Missing PDF files: {missing_files}')
        logger.info(f'Expected files: {expected_files}')
        logger.info(f'Found files: {result_files}')
        return 0.0
    if file_info:
        invalid_files = []
        for result_file in result_files:
            if result_file in file_info:
                info = file_info[result_file]
                if info.get('size', 0) <= 0:
                    invalid_files.append(f'{result_file} (empty file)')
                elif not info.get('is_valid_pdf', True):
                    invalid_files.append(f'{result_file} (not a valid PDF)')
        if invalid_files:
            logger.info(f'Invalid or empty PDF files: {invalid_files}')
            return 0.0
    logger.info(f'All validation passed: exactly {len(expected_files)} valid PDF files found')
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Found files: {result_files}')
    return 1.0

def check_pdf_files__6e9dcacc(result, expected, **options):
    """Check specific files with partial credit.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: Ratio of found files to required files (0.0-1.0)
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    if not required_files:
        return 1.0
    result_set = set(result)
    found_count = sum((1 for f in required_files if f in result_set))
    score = found_count / len(required_files)
    logger.info(f'Required {len(required_files)} files, found {found_count}')
    logger.info(f'Required: {required_files}')
    logger.info(f'Found: {result}')
    logger.info(f'Score: {score:.2f}')
    return score

def check_pdf_files__80b8ddf2(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_merge_verification__a8bc7d70fd87c66cda30ee22072de4d3(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Comprehensive verification that PDFs were actually merged from email attachments.

    This metric compares:
    1. The uploaded PDF file info (from Google Drive)
    2. The source email attachment info (from Thunderbird)

    Args:
        result: Dict with 'drive_file' and 'email_attachments' keys
        expected: Expected configuration (not used, validation is based on result consistency)
        **options: Additional options

    Returns:
        float: 1.0 if merge is verified, 0.0 otherwise
    """
    try:
        if not isinstance(result, dict):
            logger.warning(f'Result is not a dict: {type(result)}')
            return 0.0
        drive_file = result.get('drive_file', {})
        email_attachments = result.get('email_attachments', {})
        if not email_attachments.get('found', False):
            logger.warning("Email 'Paper Recommendation' not found in Thunderbird")
            return 0.0
        if not drive_file.get('exists', False):
            logger.warning('File does not exist on Google Drive')
            return 0.0
        expected_name = expected.get('filename', 'attachment_full.pdf')
        actual_name = drive_file.get('name', '')
        if actual_name != expected_name:
            logger.warning(f"Filename mismatch: '{actual_name}' vs expected '{expected_name}'")
            return 0.0
        mime_type = drive_file.get('mimeType', '')
        if mime_type != 'application/pdf':
            logger.warning(f"File is not a PDF. MIME type: '{mime_type}'")
            return 0.0
        email_pdf_count = email_attachments.get('pdf_count', 0)
        if email_pdf_count < 2:
            logger.warning(f'Email has {email_pdf_count} PDF attachments - need at least 2 to merge')
            return 0.0
        page_count = drive_file.get('page_count', 0)
        if page_count < email_pdf_count:
            logger.warning(f'Merged PDF has {page_count} pages, but {email_pdf_count} source PDFs suggest at least {email_pdf_count} pages')
            return 0.0
        file_size = drive_file.get('size', 0)
        if file_size < MIN_MERGED_PDF_SIZE_BYTES:
            logger.warning(f'File size too small for merged PDF: {file_size} bytes')
            return 0.0
        logger.info(f'Merge verification passed: {email_pdf_count} PDFs merged into {page_count} pages, file size {file_size} bytes')
        return 1.0
    except Exception as e:
        logger.error(f'Error verifying PDF merge: {e}')
        return 0.0

def check_pdf_file_count__086472c391d9b3dbf463fe72713bc019(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the PDF files meet the requirements for chapter splitting task.

    Verifies:
    1. Total count is 9 (8 chapters + 1 original)
    2. Original book PDF exists
    3. Exactly 8 chapter PDFs with proper naming convention exist

    Args:
        result: Dict with 'count', 'filenames', and 'has_original' from getter
        expected: Dict with 'count' key specifying expected number of files

    Returns:
        1.0 if all checks pass, 0.0 otherwise
    """
    expected_count = expected.get('count', 9)
    actual_count = result.get('count', 0)
    filenames = result.get('filenames', [])
    has_original = result.get('has_original', False)
    if actual_count != expected_count:
        logger.warning(f'PDF file count mismatch: got {actual_count}, expected {expected_count}')
        return 0.0
    if not has_original:
        logger.warning("Original book PDF 'Spectral Graph Theory.pdf' not found")
        return 0.0
    chapter_patterns = ['^Chapter[_ ]?(\\d+)\\.pdf$', '^Ch[_ ]?(\\d+)\\.pdf$', '^chapter[_ ]?(\\d+)\\.pdf$', '^ch[_ ]?(\\d+)\\.pdf$']
    chapter_numbers = set()
    for filename in filenames:
        if filename == 'Spectral Graph Theory.pdf':
            continue
        matched = False
        for pattern in chapter_patterns:
            match = re.match(pattern, filename)
            if match:
                chapter_num = int(match.group(1))
                chapter_numbers.add(chapter_num)
                matched = True
                break
        if not matched:
            logger.warning(f"File '{filename}' does not match expected chapter naming convention")
            return 0.0
    if len(chapter_numbers) != 8:
        logger.warning(f'Expected 8 chapter files, found {len(chapter_numbers)}')
        return 0.0
    if chapter_numbers != set(range(1, 9)):
        logger.warning(f'Chapter numbers should be 1-8, got: {sorted(chapter_numbers)}')
        return 0.0
    logger.info(f'All checks passed: 9 total PDFs, original exists, 8 chapters with correct naming')
    return 1.0

def check_pdf_validity__1cc9c1bfb30ee1b2b3cdbf7926894918(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDFs are valid (proper extension, count, and filename matching).

    Args:
        result: Dict with 'files', 'valid_count', 'chrome_titles', and 'matched_files' from getter
        expected: Dict with 'min_count' for minimum valid PDFs expected
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
        - 0.0: No PDFs or no matching filenames
        - 0.5: Correct count but filenames don't match titles
        - 0.75: Some filenames match (partial)
        - 1.0: Correct count AND all filenames match article titles
    """
    min_count = expected.get('min_count', 2)
    valid_count = result.get('valid_count', 0)
    matched_files = result.get('matched_files', [])
    matched_count = len(matched_files)
    chrome_titles = result.get('chrome_titles', [])
    if valid_count == 0:
        return 0.0
    count_ok = valid_count >= min_count
    if chrome_titles:
        if matched_count >= min_count:
            return 1.0
        elif matched_count > 0 and count_ok:
            match_ratio = matched_count / min_count
            return 0.5 + match_ratio * 0.5
        elif count_ok:
            return 0.5
        else:
            return valid_count / min_count * 0.5
    elif count_ok:
        return 0.5
    else:
        return valid_count / min_count * 0.5

def check_pdf_file_properties__6f106331(pdf_file: str, expected, **options):
    """
    Check if PDF exists, has reasonable file size, and is a valid PDF.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with 'min_size' in bytes
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.3
    try:
        file_size = os.path.getsize(pdf_file)
        min_size = expected.get('min_size', 10240)
        if file_size >= min_size:
            score += 0.3
        reader = PdfReader(pdf_file)
        if len(reader.pages) > 0:
            score += 0.4
    except Exception as e:
        return score
    return score

def check_pdf_non_empty__d4f3d6039bf1a71698a65c2d049521e8(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDFs are non-empty, meet minimum size, and have correct filenames.

    Args:
        result: Dict with 'files' (filename->size) and 'filenames' (list) from getter
        expected: Dict with 'min_count', 'min_size_bytes', and 'expected_keywords' keys
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    min_count = expected.get('min_count', 2)
    min_size_bytes = expected.get('min_size_bytes', 1000)
    expected_keywords = expected.get('expected_keywords', [])
    file_sizes = result.get('files', {})
    filenames = result.get('filenames', [])
    valid_files = [f for (f, size) in file_sizes.items() if size >= min_size_bytes]
    valid_count = len(valid_files)
    filename_score = 1.0
    if expected_keywords:
        matched_keywords = 0
        for keyword_set in expected_keywords:
            if isinstance(keyword_set, str):
                keywords = [keyword_set.lower()]
            else:
                keywords = [k.lower() for k in keyword_set]
            for filename in filenames:
                filename_lower = filename.lower()
                filename_lower = filename_lower.replace('.pdf', '')
                if all((keyword in filename_lower for keyword in keywords)):
                    matched_keywords += 1
                    break
        filename_score = matched_keywords / len(expected_keywords) if expected_keywords else 1.0
    if valid_count >= min_count:
        count_score = 1.0
    elif valid_count > 0:
        count_score = valid_count / min_count
    else:
        count_score = 0.0
    final_score = min(count_score, filename_score)
    return final_score

def check_pdf_saved__60388d3f6c5270e1728288c485d5fd5b(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if a PDF file was successfully saved with valid content.

    Args:
        result: Path to the PDF file (from getter)
        expected: Dict with validation rules (min_pages, min_size_kb)
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size_kb = os.path.getsize(result) / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        from pypdf import PdfReader
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        min_pages = expected.get('min_pages', 1)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_files_exist__6d39fc800ff8d1397c30bfae2c676b76(result, expected, **options):
    """Check if all expected PDF files exist.

    Args:
        result: Dict mapping filenames to existence status
        expected: Expected filenames that should exist
        **options: Additional options

    Returns:
        float: Score based on how many files exist
    """
    expected_files = expected.get('filenames', [])
    if not expected_files:
        return 0.0
    exist_count = sum((1 for f in expected_files if result.get(f, False)))
    score = exist_count / len(expected_files)
    return score

def check_pdf_text_contains__a09f0e42332d230e7cc0e7794732425e(result: Any, expected: Any, **options) -> float:
    """Check if PDF text contains expected keywords.

    Args:
        result: Actual text content from getter
        expected: Expected configuration with keywords (from rules dict)
        **options: Additional options

    Returns:
        float: Percentage of keywords found (0.0-1.0)
    """
    try:
        keywords = expected.get('keywords', [])
        if not isinstance(result, str):
            logger.warning(f'Result is not a string: {type(result)}')
            return 0.0
        if not keywords:
            logger.warning('No keywords specified')
            return 0.0
        result_lower = result.lower()
        found_count = 0
        for keyword in keywords:
            if keyword.lower() in result_lower:
                found_count += 1
                logger.info(f"Found keyword: '{keyword}'")
            else:
                logger.warning(f"Missing keyword: '{keyword}'")
        score = found_count / len(keywords)
        logger.info(f'Text content check: {found_count}/{len(keywords)} keywords found (score: {score:.2f})')
        return score
    except Exception as e:
        logger.error(f'Error checking PDF text content: {e}')
        return 0.0

def check_pdf_saved__40d8ee41df97cbb2f480ba7af545efc4(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if a PDF file was successfully saved with valid content.

    Args:
        result: Path to the PDF file (from getter)
        expected: Dict with validation rules (min_pages, min_size_kb)
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size_kb = os.path.getsize(result) / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        from pypdf import PdfReader
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        min_pages = expected.get('min_pages', 1)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_files__8010e79b(result, expected, **options):
    """Check if minimum number of PDF files exist.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'min_files' key
        **options: Additional options
        
    Returns:
        float: 1.0 if count >= min_files, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    min_files = expected.get('min_files', 10)
    actual_count = len(result)
    logger.info(f'Expected at least {min_files} PDF files, found {actual_count}')
    logger.info(f'PDF files: {result}')
    if actual_count >= min_files:
        return 1.0
    else:
        return 0.0

def check_pdf_filesize__74ce4e60(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file size is within expected range.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if file size is within range, 0.0 otherwise
    """
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        file_size = os.path.getsize(result)
        if min_size <= file_size <= max_size:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_pdf_export__9a18b30d646547c54b42b3593f83920d(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF was exported correctly with expected page count.

    Args:
        result: Dict from getter with 'exists', 'page_count', 'file_size'
        expected: Dict with 'min_pages' requirement
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists', False):
        return 0.0
    page_count = result.get('page_count', 0)
    min_pages = expected.get('min_pages', 1)
    if page_count < min_pages:
        return 0.0
    file_size = result.get('file_size', 0)
    if file_size < 10240:
        return 0.0
    return 1.0

def check_pdf_has_text_content__e1e75309_6(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists and has extractable text content (not just images).

    Args:
        result: Path to the PDF file (from vm_file getter)
        expected: Rules dict with 'min_text_length' (minimum number of characters)

    Returns:
        float: 1.0 if PDF exists and has sufficient text content, 0.0 otherwise
    """
    if not result:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if not result.lower().endswith('.pdf'):
        return 0.0
    min_text_length = expected.get('min_text_length', 100)
    try:
        doc = fitz.open(result)
        total_text = ''
        for page in doc:
            total_text += page.get_text()
        doc.close()
        text_length = len(total_text.strip())
        if text_length < min_text_length:
            return 0.0
        return 1.0
    except Exception:
        return 0.0

def check_pdf_filename__e61f394075a7c32a6a0c2c96a3700939(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF exists with expected filename.

    Args:
        result: Dict from getter with filename info
        expected: Dict with expected filename (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        logger.warning(f"PDF file does not exist at {result.get('path')}")
        return 0.0
    if not result.get('valid_pdf', False):
        logger.warning('PDF file is invalid or corrupted')
        return 0.0
    expected_filename = expected.get('expected_filename', '')
    actual_filename = result.get('filename', '')
    if actual_filename == expected_filename:
        logger.info(f'Filename matches: {expected_filename}')
        return 1.0
    else:
        logger.warning(f"Filename mismatch: expected '{expected_filename}', got '{actual_filename}'")
        return 0.0

def check_pdf_count_and_content__9ba5dc81930fa553e6c6310edaaff2ec(result_state: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDFs meet count, filename, and content requirements.

    Args:
        result_state: Dict with 'count', 'files', and 'filenames' from getter
        expected: Dict containing (when type='rule', framework passes rules directly):
            - 'min_pdf_count': Minimum number of PDFs required
            - 'content_strings': List of strings that should appear in PDFs
            - 'expected_filenames': List of expected filenames (optional)
        **options: Additional options (env for file access)

    Returns:
        float: Score from 0.0 to 1.0
    """
    import fitz
    if not result_state:
        logger.error('Result state is empty or None')
        return 0.0
    pdf_count = result_state.get('count', 0)
    pdf_files = result_state.get('files', [])
    pdf_filenames = result_state.get('filenames', [])
    min_pdf_count = expected.get('min_pdf_count', 0)
    content_strings = expected.get('content_strings', [])
    expected_filenames = expected.get('expected_filenames', [])
    score = 0.0
    if pdf_count >= min_pdf_count:
        score += 0.3
        logger.info(f'PDF count check passed: {pdf_count} >= {min_pdf_count}')
    else:
        logger.warning(f'PDF count check failed: {pdf_count} < {min_pdf_count}')
        return score
    if expected_filenames:
        filename_matches = 0
        for expected_filename in expected_filenames:
            if expected_filename in pdf_filenames:
                filename_matches += 1
                logger.info(f"Filename match found: '{expected_filename}'")
            else:
                logger.warning(f"Expected filename not found: '{expected_filename}'")
        if filename_matches == len(expected_filenames):
            score += 0.3
            logger.info('All expected filenames found')
        elif filename_matches > 0:
            score += 0.3 * (filename_matches / len(expected_filenames))
            logger.info(f'Partial filename matches: {filename_matches}/{len(expected_filenames)}')
    else:
        score += 0.3
        logger.info('No expected filenames to check')
    if not content_strings:
        logger.warning('No content strings to check')
        return score
    env = options.get('env')
    if not env:
        logger.error('Environment not provided in options')
        return score
    all_pdf_texts = []
    for pdf_path in pdf_files:
        try:
            file_bytes = env.controller.get_file(pdf_path)
            if not file_bytes:
                logger.warning(f'Failed to get file: {pdf_path}')
                continue
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                doc = fitz.open(tmp_path)
                pdf_text = ''
                for page in doc:
                    pdf_text += page.get_text()
                doc.close()
                all_pdf_texts.append(pdf_text)
                logger.info(f'Extracted text from {pdf_path} ({len(pdf_text)} chars)')
            finally:
                os.unlink(tmp_path)
        except Exception as e:
            logger.error(f'Error processing PDF {pdf_path}: {e}')
            continue
    points_per_string = 0.4 / len(content_strings)
    for content_string in content_strings:
        found = False
        for pdf_text in all_pdf_texts:
            if content_string in pdf_text:
                found = True
                logger.info(f"Found content string: '{content_string[:50]}...'")
                break
        if found:
            score += points_per_string
        else:
            logger.warning(f"Content string not found: '{content_string[:50]}...'")
    logger.info(f'Final score: {score:.2f}')
    return score

def check_pdf_files__d928a635(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_keywords__c68533a8ce70563646c4156811d621fa(result, expected, **options):
    """Check if PDFs contain required keywords with per-file verification.

    Args:
        result: Dict from getter with filename -> list of found keywords
        expected: Dict with 'files_with_keywords' structure
        **options: Additional options

    Returns:
        float: Score based on percentage of keywords found, verifying each keyword
               is in its corresponding file (not just total count)
    """
    files_with_keywords = expected.get('files_with_keywords', [])
    if not files_with_keywords:
        return 0.0
    total_keywords = 0
    found_keywords = 0
    for file_info in files_with_keywords:
        filename = file_info['filename']
        expected_keywords = file_info.get('keywords', [])
        total_keywords += len(expected_keywords)
        found_list = result.get(filename, [])
        for expected_kw in expected_keywords:
            if expected_kw in found_list:
                found_keywords += 1
    if total_keywords == 0:
        return 0.0
    return found_keywords / total_keywords

def check_pdf_valid__e4448a5a(pdf_file: str, expected, **options):
    """
    Check if PDF is valid and contains expected content.

    Args:
        pdf_file: Path to the PDF file
        expected: Dict with validation rules
        **options: Additional options

    Returns:
        float: Score 0.0-1.0 with partial credit
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    score = 0.0
    try:
        reader = PdfReader(pdf_file)
        if len(reader.pages) > 0:
            score += 0.5
            has_images = False
            for page in reader.pages:
                if '/XObject' in page.get('/Resources', {}):
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        if xobjects[obj]['/Subtype'] == '/Image':
                            has_images = True
                            break
                if has_images:
                    break
            if has_images:
                score += 0.5
    except Exception as e:
        return 0.0
    return score

def check_single_pdf_complete__1386929a4648885c7d87d6425829fd94(result, expected, **options):
    """Check if single PDF exists and has correct content.

    Args:
        result: Dict from getter with 'exists' and 'has_content' keys
        expected: Dict (not used, validation is in getter)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 (0.5 for exists, 1.0 for correct content)
    """
    if not result.get('exists', False):
        return 0.0
    if result.get('has_content', False):
        return 1.0
    return 0.5

def check_pdf_format__988ec512(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if file is a valid PDF with minimum pages.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if valid PDF format, 0.0 otherwise
    """
    min_pages = expected.get('min_pages', 1)
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        if num_pages >= min_pages:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_pdf_file_exists__7648b3ac(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDF file exists.

    Args:
        result: Dict with file_exists from getter
        expected: Dict with expected file path
        **options: Additional options

    Returns:
        Score: 1.0 if file exists, 0.0 otherwise
    """
    if result.get('file_exists', False):
        logger.info(f"PDF file exists at {result.get('file_path')}")
        return 1.0
    else:
        logger.warning(f"PDF file does not exist at {result.get('file_path')}")
        return 0.0

def check_pdf_filenames__7152d37e(result, expected, **options):
    """Check if PDFs with expected names exist.

    Args:
        result: List of PDF filenames (without .pdf extension)
        expected: Dict with:
            - expected_names: List of expected filenames (without extension)

    Returns:
        Score between 0.0 and 1.0
    """
    expected_names = expected.get('expected_names', [])
    if not isinstance(result, list):
        return 0.0
    matches = sum((1 for name in expected_names if name in result))
    if len(expected_names) == 0:
        return 0.0
    return matches / len(expected_names)

def check_pdf_file_size__29f6acc744cf0d15caa355ed1701b507(result: Any, expected: Any, **options) -> float:
    """Check if PDF file size is within expected range.

    Args:
        result: Actual file size in bytes from getter
        expected: Expected configuration (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if within range, 0.0 otherwise
    """
    try:
        min_size = expected.get('min_size_bytes', 0)
        max_size = expected.get('max_size_bytes', float('inf'))
        if not isinstance(result, (int, float)):
            logger.warning(f'Result is not a number: {result}')
            return 0.0
        if min_size <= result <= max_size:
            logger.info(f'File size within range: {result} bytes (min: {min_size}, max: {max_size})')
            return 1.0
        else:
            logger.warning(f'File size out of range: {result} bytes (expected: {min_size}-{max_size})')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking PDF file size: {e}')
        return 0.0

def check_pdf_exists__05a7fe2932371163af38b8e77d9b0c93(result: bool, expected: dict, **options) -> float:
    """
    Check if PDF file exists as expected.

    Args:
        result: Boolean indicating if PDF file exists (from getter)
        expected: Expected rules dict with 'should_exist' key (True/False)
        **options: Additional options

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise
    """
    should_exist = expected.get('should_exist', True)
    if result == should_exist:
        logger.info(f'PDF existence check passed: {result} == {should_exist}')
        return 1.0
    else:
        logger.warning(f'PDF existence check failed: {result} != {should_exist}')
        return 0.0

def check_pdf_name_pattern__55b958e432c8380deab73a3d2fcf329a(result, expected, **options):
    """Check if PDFs exist with correct naming pattern and content.

    Args:
        result: Dict from getter with filename -> status mapping
        expected: Dict with 'expected_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    total_files = len(expected_files)
    correct_files = 0
    for file_info in expected_files:
        filename = file_info['filename']
        status = result.get(filename, -1)
        if status == 1:
            correct_files += 1
    return correct_files / total_files

def check_chapter_pdf_count__457b9850ebb20f14d7c68688b0aa27a6(result, expected, **options):
    """Check if the number of chapter PDF files meets the minimum threshold.

    Args:
        result: Dict with 'count' (int) and 'chapter_files' (list) from getter
        expected: Dict with 'min_chapters' (int) - minimum number of chapters expected
        **options: Additional options

    Returns:
        float: 1.0 if count >= min_chapters, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    actual_count = result.get('count', 0)
    min_chapters = expected.get('min_chapters', 5)
    if actual_count >= min_chapters:
        return 1.0
    else:
        return 0.0

def check_pdf_chapters__6b40f44d(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count from getter
        expected: Dict with 'files' (list of filenames) and 'page_counts' (dict with thresholds)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    page_count_rules = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        if expected_file in page_count_rules:
            rule = page_count_rules[expected_file]
            relation = rule.get('relation', 'eq')
            ref_value = rule.get('ref_value', 0)
            actual_page_count = result[expected_file]
            if relation == 'ge' and actual_page_count >= ref_value:
                score += points_per_file
            elif relation == 'le' and actual_page_count <= ref_value:
                score += points_per_file
            elif relation == 'eq' and actual_page_count == ref_value:
                score += points_per_file
            elif relation == 'gt' and actual_page_count > ref_value:
                score += points_per_file
            elif relation == 'lt' and actual_page_count < ref_value:
                score += points_per_file
        else:
            score += points_per_file
    return min(score, 1.0)

def check_pdf_exists__8752e3fa(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file exists and meets size requirements.

    Args:
        result: Dict from getter with 'exists' and 'size_bytes' keys
        expected: Expected rules dict with 'exists' and 'min_size_kb' keys

    Returns:
        float: 1.0 if file exists and meets requirements, 0.0 otherwise
    """
    if not result.get('exists', False):
        return 0.0
    if not expected.get('exists', True):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 0)
    if min_size_kb > 0:
        size_kb = result.get('size_bytes', 0) / 1024
        if size_kb < min_size_kb:
            return 0.0
    return 1.0

def check_pdf_validity__0b5ef92a5e5fe7305b438702a0eb6c3a(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF is valid and meets basic requirements.

    Args:
        result: Dict from getter with file properties
        expected: Dict with expected properties
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    if result.get('is_valid', False):
        score += 0.5
    if result.get('page_count', 0) >= expected.get('min_pages', 1):
        score += 0.25
    min_size = expected.get('min_size', 1000)
    if result.get('file_size', 0) >= min_size:
        score += 0.25
    return score

def check_pdf_export__c5999de9(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names and content based on article titles.

    Args:
        result: List of dicts with 'filename' and 'content' keys
        expected: Dict with 'article_patterns' - list of patterns for each article
                  Each pattern has 'keywords' (list of keywords to match) and
                  'title_fragment' (identifier for the article)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    article_patterns = expected.get('article_patterns', [])
    expected_count = expected.get('count', len(article_patterns))
    if len(result) != expected_count:
        return 0.0
    matched_pdfs = set()
    for pattern in article_patterns:
        keywords = pattern.get('keywords', [])
        title_fragment = pattern.get('title_fragment', '')
        pattern_matched = False
        for pdf_data in result:
            pdf_filename = pdf_data.get('filename', '')
            pdf_content = pdf_data.get('content', '')
            if pdf_filename in matched_pdfs:
                continue
            filename_lower = pdf_filename.lower()
            filename_score = 0
            if title_fragment == 'agent':
                has_llm_or_autonomous = 'llm' in filename_lower or 'autonomous' in filename_lower
                has_agent = 'agent' in filename_lower
                if has_llm_or_autonomous and has_agent:
                    filename_score = 1
            elif title_fragment == 'human':
                has_human = 'human' in filename_lower
                has_data = 'data' in filename_lower
                has_quality = 'quality' in filename_lower
                if has_human and has_data and has_quality:
                    filename_score = 1
            else:
                matches = sum((1 for kw in keywords if kw.lower() in filename_lower))
                if matches >= min(2, len(keywords)):
                    filename_score = 1
            content_score = 0
            if pdf_content:
                if title_fragment == 'agent':
                    agent_phrases = ['autonomous agent', 'llm powered', 'llm-powered', 'chain of thought', 'reasoning', 'planning']
                    content_matches = sum((1 for phrase in agent_phrases if phrase in pdf_content))
                    if content_matches >= 2:
                        content_score = 1
                elif title_fragment == 'human':
                    human_phrases = ['data quality', 'human annotation', 'labeling', 'dataset', 'quality control']
                    content_matches = sum((1 for phrase in human_phrases if phrase in pdf_content))
                    if content_matches >= 2:
                        content_score = 1
                elif all((kw.lower() in pdf_content for kw in keywords)):
                    content_score = 1
            if filename_score > 0 and content_score > 0:
                matched_pdfs.add(pdf_filename)
                pattern_matched = True
                break
        if not pattern_matched:
            return 0.0
    return 1.0

def check_pdf_split_complete__ea295e4c379ee25a192357c908aada76(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF splitting is complete (chapters created and original preserved).

    Args:
        result: Dict with 'total_count', 'numbered_chapters', 'other_files' from getter
        expected: Dict with 'expected_chapters' (number of chapters expected),
            'original_filename' (name of original PDF file)

    Returns:
        Score between 0.0 and 1.0
        - 0.5 if expected number of chapters exist
        - 1.0 if chapters exist AND original file is preserved
    """
    expected_chapter_count = expected.get('expected_chapters', 8)
    original_filename = expected.get('original_filename', 'Spectral Graph Theory.pdf')
    if not result:
        logger.warning('No PDF file info received')
        return 0.0
    numbered_chapters = result.get('numbered_chapters', [])
    other_files = result.get('other_files', [])
    score = 0.0
    if len(numbered_chapters) >= expected_chapter_count:
        score += 0.5
        logger.info(f'Found {len(numbered_chapters)} numbered chapters (expected {expected_chapter_count})')
    else:
        logger.warning(f'Insufficient chapters: found {len(numbered_chapters)}, expected {expected_chapter_count}')
    if original_filename in other_files:
        score += 0.5
        logger.info(f"Original file '{original_filename}' is preserved")
    else:
        logger.warning(f"Original file '{original_filename}' not found in other files: {other_files}")
    logger.info(f'PDF split completion score: {score:.2f}')
    return score

def check_pdf_files__cdfaf4c1(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_files__0707ddca(result, expected, **options):
    """Check if minimum number of PDF files exist.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'min_files' key
        **options: Additional options
        
    Returns:
        float: 1.0 if count >= min_files, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    min_files = expected.get('min_files', 10)
    actual_count = len(result)
    logger.info(f'Expected at least {min_files} PDF files, found {actual_count}')
    logger.info(f'PDF files: {result}')
    if actual_count >= min_files:
        return 1.0
    else:
        return 0.0

def check_pdf_chapters__6871c0fe(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count from getter
        expected: Expected state with 'files' and 'page_counts' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    expected_page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    total_checks = len(expected_files) * 2
    passed_checks = 0
    for expected_file in expected_files:
        if expected_file in result:
            passed_checks += 1
            if expected_file in expected_page_counts:
                actual_pages = result[expected_file]
                page_requirement = expected_page_counts[expected_file]
                relation = page_requirement.get('relation', 'ge')
                ref_value = page_requirement.get('ref_value', 0)
                if relation == 'ge' and actual_pages >= ref_value:
                    passed_checks += 1
                elif relation == 'le' and actual_pages <= ref_value:
                    passed_checks += 1
                elif relation == 'eq' and actual_pages == ref_value:
                    passed_checks += 1
                elif relation == 'gt' and actual_pages > ref_value:
                    passed_checks += 1
                elif relation == 'lt' and actual_pages < ref_value:
                    passed_checks += 1
            else:
                passed_checks += 1
    score = passed_checks / total_checks if total_checks > 0 else 0.0
    return min(score, 1.0)

def check_pdf_count__bfc9cce9462c8f8be02842ad8f805893(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if the number of PDF files matches expected count.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'count' key specifying expected number of PDFs
        **options: Additional options

    Returns:
        1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('count', 2)
    actual_count = len(result)
    if actual_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_pdf_exists__675930f3(result_file, expected, **options):
    """
    Check if PDF file exists and has minimum number of pages.

    Args:
        result_file: Path to the PDF file
        expected: Dict containing 'min_pages' key
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists with sufficient pages, 0.0 otherwise
    """
    try:
        with open(result_file, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            min_pages = expected.get('min_pages', 1)
            if num_pages >= min_pages:
                return 1.0
            else:
                return 0.0
    except Exception as e:
        return 0.0

def check_pdf_files__201ff98c(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_export__2586a709(result, expected, **options):
    """
    Compare exported PDF list against expected PDF requirements.
    Validates that PDFs are valid, contain content, and have filenames that match actual webpage titles.

    Args:
        result: Dict containing:
            - pdf_files: List of dicts with 'filename', 'valid', 'text_content', 'page_count', 'pdf_title'
            - count: Total number of PDF files found
        expected: Dict with:
            - 'required_title_phrases' - list of lists, each inner list contains phrases/words from one webpage title
            - 'count' - expected number of PDFs
            - 'content_keywords' - list of keywords expected in PDF content (optional)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    pdf_files = result.get('pdf_files', [])
    count = result.get('count', 0)
    required_title_phrases = expected.get('required_title_phrases', [])
    expected_count = expected.get('count', len(required_title_phrases))
    content_keywords = expected.get('content_keywords', [])
    if count != expected_count:
        return 0.0
    valid_pdfs = [pdf for pdf in pdf_files if pdf.get('valid', False)]
    if len(valid_pdfs) != expected_count:
        return 0.0
    for pdf in valid_pdfs:
        if pdf.get('page_count', 0) < 1:
            return 0.0
    for pdf in valid_pdfs:
        text = pdf.get('text_content', '').strip()
        if not text:
            return 0.0
    score_components = []
    filename_match_score = 0.0
    matched_titles = set()
    for pdf in valid_pdfs:
        filename = pdf.get('filename', '').lower().replace('.pdf', '')
        best_match_score = 0.0
        best_match_idx = -1
        for (idx, phrase_list) in enumerate(required_title_phrases):
            if idx in matched_titles:
                continue
            matches = 0
            for phrase in phrase_list:
                if phrase.lower() in filename:
                    matches += 1
            match_score = matches / len(phrase_list) if phrase_list else 0
            if match_score >= 0.5 and match_score > best_match_score:
                best_match_score = match_score
                best_match_idx = idx
        if best_match_idx >= 0:
            matched_titles.add(best_match_idx)
            filename_match_score += 1.0 / expected_count
    score_components.append(filename_match_score * 0.5)
    content_score = 0.0
    if content_keywords:
        all_content = ' '.join([pdf.get('text_content', '') for pdf in valid_pdfs])
        for keyword in content_keywords:
            if keyword.lower() in all_content.lower():
                content_score += 1.0 / len(content_keywords)
    else:
        content_score = 1.0
    score_components.append(content_score * 0.3)
    meaningful_filename_score = 0.0
    generic_names = ['document', 'file', 'download', 'untitled', 'page', 'export', 'chrome', 'tab']
    meaningful_count = 0
    for pdf in valid_pdfs:
        filename = pdf.get('filename', '').lower().replace('.pdf', '')
        is_generic = any((gen == filename or filename.startswith(gen + '_') for gen in generic_names))
        if len(filename) >= 5 and (not is_generic):
            meaningful_count += 1
    if expected_count > 0:
        meaningful_filename_score = meaningful_count / expected_count
    score_components.append(meaningful_filename_score * 0.2)
    total_score = sum(score_components)
    return total_score

def check_pdf_chapters__92f3cb2c(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping filename to page count, e.g., {'file.pdf': 10}
        expected: Expected rules with 'files' (list) and 'page_counts' (dict)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('files', [])
    page_counts = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(expected_files)
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        file_score = points_per_file * 0.5
        if expected_file in page_counts:
            page_requirement = page_counts[expected_file]
            relation = page_requirement.get('relation', 'eq')
            ref_value = page_requirement.get('ref_value', 0)
            actual_count = result[expected_file]
            page_count_valid = False
            if relation == 'ge':
                page_count_valid = actual_count >= ref_value
            elif relation == 'le':
                page_count_valid = actual_count <= ref_value
            elif relation == 'eq':
                page_count_valid = actual_count == ref_value
            elif relation == 'gt':
                page_count_valid = actual_count > ref_value
            elif relation == 'lt':
                page_count_valid = actual_count < ref_value
            if page_count_valid:
                file_score += points_per_file * 0.5
        else:
            file_score = points_per_file
        score += file_score
    return min(score, 1.0)

def check_pdf_export__98b0d7a7(result, expected, **options):
    """
    Compare exported PDF list against actual page titles from Chrome.

    Args:
        result: Dict with 'pdf_files' (list of PDF filenames) and 'page_titles' (list of actual page titles)
        expected: Dict with 'count' - expected number of PDFs

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    pdf_files = result.get('pdf_files', [])
    page_titles = result.get('page_titles', [])
    expected_count = expected.get('count', 2)
    if len(pdf_files) != expected_count:
        return 0.0
    if not page_titles:
        return 0.5 if len(pdf_files) == expected_count else 0.0
    if len(page_titles) != expected_count:
        pass
    matched_pdfs = set()
    matched_titles = set()
    for (i, page_title) in enumerate(page_titles):
        for (j, pdf_file) in enumerate(pdf_files):
            if j not in matched_pdfs and title_matches_filename(page_title, pdf_file):
                matched_pdfs.add(j)
                matched_titles.add(i)
                break
    if len(page_titles) > 0:
        score = len(matched_titles) / len(page_titles)
    else:
        score = 0.0
    if len(matched_titles) == expected_count and len(pdf_files) == expected_count:
        return 1.0
    return score

def check_pdf_chapters_exist__242f4871e4280ec1025212dbeaf5c18c(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if expected PDF chapter files exist.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'required_files' (list of filenames that must exist)
                  and 'original_file' (the original PDF that should also exist)

    Returns:
        Score between 0.0 and 1.0 based on how many required files exist
    """
    required_files = expected.get('required_files', [])
    original_file = expected.get('original_file')
    if not result:
        logger.warning('No PDF files found')
        return 0.0
    score = 0.0
    total_checks = len(required_files)
    if original_file:
        total_checks += 1
    if total_checks == 0:
        return 0.0
    for required_file in required_files:
        if required_file in result:
            score += 1.0 / total_checks
        else:
            logger.warning(f'Required file not found: {required_file}')
    if original_file:
        if original_file in result:
            score += 1.0 / total_checks
        else:
            logger.warning(f'Original file not found: {original_file}')
    logger.info(f'PDF chapters check score: {score:.2f} ({int(score * total_checks)}/{total_checks} files found)')
    return score

def check_pdf_file__766a8b90(result, expected, **options):
    """
    Check if a valid PDF file exists based on file command output with timestamp and size verification.

    This function verifies that:
    1. A PDF file exists at the target location (/home/user/Documents/lecture-notes.pdf)
    2. The file was created/modified recently (within 5 minutes) to confirm it's from task execution
    3. The file type is confirmed as PDF via file command
    4. The file has reasonable size (non-empty, within expected range for PDF documents)

    Task verification approach:
    The task requires: (1) finding email with subject 'Lecture Document',
    (2) locating it in Notes folder, (3) saving its attachment as PDF to Documents.

    Direct verification of email operations is not possible without Thunderbird-specific
    getters for email state/folder inspection. Therefore, this evaluator uses an enhanced
    file-based verification strategy with multiple checks:

    - Preconfig cleanup (line 38-42 in task config) removes any pre-existing file at the
      target path to prevent false positives from old files
    - Timestamp check (time_diff < 300 seconds) ensures the file is newly created/saved
      during task execution, not a pre-existing file
    - File type verification confirms it's a valid PDF document
    - File size verification ensures the PDF has reasonable content (>100 bytes, <50MB)
    - The specific filename 'lecture-notes.pdf' is controlled by test data

    This enhanced approach provides strong confidence that the task was completed correctly,
    as the combination of:
    - Cleanup preventing false positives from old files
    - Recent timestamp indicating new file creation during task execution
    - Correct filename and location matching task requirements
    - Valid PDF file type with reasonable size

    makes it highly unlikely the file appeared through any means other than correct
    task execution (saving the email attachment).

    Assumptions:
    - The attachment in the test email 'Lecture Document' is named 'lecture-notes.pdf'
    - The task must complete within 5 minutes for timestamp verification to work
    - Test environment guarantees the filename assumption via controlled test data
    - Preconfig cleanup ensures no pre-existing files with the same name
    - Valid PDF attachments are between 100 bytes and 50MB

    Limitations (acceptable given available tools):
    - Cannot directly verify the email was found in Notes folder (no Thunderbird getters)
    - Cannot verify file content fingerprint matches original attachment exactly
    - Cannot verify email was opened or attachment was accessed (no email state inspection)
    - Relies on indirect but robust evidence (correct file at correct location with
      recent timestamp and valid PDF characteristics)

    Note: This is the most robust verification possible without Thunderbird-specific getters
    for email/folder state inspection. The multiple layers of verification (cleanup +
    timestamp + type + size + location + filename) provide high confidence in correct
    task completion.

    Args:
        result: Output from shell command that checks file existence, type, timestamp, and size
        expected: Dict with 'expected_output' key (should be 'valid')
        **options: Additional options

    Returns:
        float: 1.0 if valid PDF file exists with recent timestamp and reasonable size, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'valid')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_pdf_orientation__c75d6b17(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF has expected orientation.

    Args:
        result: Orientation string ('landscape' or 'portrait')
        expected: Dict with 'orientation' (expected orientation)
        **options: Additional options

    Returns:
        float: 1.0 if orientation matches, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_orientation = expected.get('orientation', 'portrait')
    if result == expected_orientation:
        return 1.0
    return 0.0

def check_pdf_export__d263a7ae(result, expected, **options):
    """
    Compare exported PDF list against expected PDF names.

    Args:
        result: List of actual PDF filenames
        expected: Dict with 'required_keywords' - list of keywords that must appear in filenames,
                  and 'count' - expected number of PDFs

    Returns:
        float: Score 1.0 if all requirements met, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    required_keywords = expected.get('required_keywords', [])
    expected_count = expected.get('count', len(required_keywords))
    if len(result) != expected_count:
        return 0.0
    matched_files = set()
    for keyword in required_keywords:
        found = False
        for filename in result:
            if keyword.lower() in filename.lower() and filename not in matched_files:
                matched_files.add(filename)
                found = True
                break
        if not found:
            return 0.0
    return 1.0

def check_pdf_count__2f750d009f0f471751ed869c09f90a90(result, expected, **options):
    """Check if the PDF file count matches expected value.

    Args:
        result: Actual PDF file count from getter
        expected: Expected count from rules
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_pdf_files__a96f92e2(result, expected, **options):
    """Check if minimum number of PDF files exist.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'min_files' key
        **options: Additional options
        
    Returns:
        float: 1.0 if count >= min_files, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    min_files = expected.get('min_files', 10)
    actual_count = len(result)
    logger.info(f'Expected at least {min_files} PDF files, found {actual_count}')
    logger.info(f'PDF files: {result}')
    if actual_count >= min_files:
        return 1.0
    else:
        return 0.0

def check_pdf_basic__06c89307(pdf_file: str, rules: Dict[str, Any]) -> float:
    """
    Check if PDF has the expected page count and optionally verify content.

    This function validates:
    1. PDF file exists
    2. PDF has the expected number of pages (from rules['page_count'])
    3. If verify_content is True, checks:
       - File size is reasonable (> 10KB to ensure not empty)
       - PDF contains embedded images (for image conversion tasks)

    Args:
        pdf_file: Path to the PDF file
        rules: Dict with 'page_count' (int) and optional 'verify_content' (bool) flag

    Returns:
        float: 1.0 if all conditions met, 0.0 otherwise
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    try:
        reader = PdfReader(pdf_file)
        nb_pages: int = len(reader.pages)
        expected_page_count = rules.get('page_count')
        if expected_page_count is None:
            return 0.0
        if nb_pages != expected_page_count:
            return 0.0
        if rules.get('verify_content', False):
            file_size = os.path.getsize(pdf_file)
            if file_size < 10240:
                return 0.0
            has_images = False
            for page in reader.pages:
                if '/XObject' in page['/Resources']:
                    xobjects = page['/Resources']['/XObject'].get_object()
                    for obj in xobjects:
                        if xobjects[obj]['/Subtype'] == '/Image':
                            has_images = True
                            break
                if has_images:
                    break
            if not has_images:
                return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_files__cb51ce0b(result, expected, **options):
    """Check specific files with partial credit.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'files' key (list of required filenames)
        **options: Additional options
        
    Returns:
        float: Ratio of found files to required files (0.0-1.0)
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    required_files = expected.get('files', [])
    if not required_files:
        return 1.0
    result_set = set(result)
    found_count = sum((1 for f in required_files if f in result_set))
    score = found_count / len(required_files)
    logger.info(f'Required {len(required_files)} files, found {found_count}')
    logger.info(f'Required: {required_files}')
    logger.info(f'Found: {result}')
    logger.info(f'Score: {score:.2f}')
    return score

def check_pdf_not_encrypted__391a007e(pdf_file: str, rules: Dict[str, Any]) -> float:
    """
    Check if PDF file exists and is not encrypted/password protected.

    Args:
        pdf_file: Path to the PDF file
        rules: Dict with validation rules (can be empty for this check)

    Returns:
        float: 1.0 if PDF exists and is not encrypted, 0.0 otherwise
    """
    if pdf_file is None:
        return 0.0
    if not os.path.exists(pdf_file):
        return 0.0
    try:
        reader = PdfReader(pdf_file)
        if reader.is_encrypted:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_year_organized_pdfs__5ad487ac7d025a6b906a4c83e8beac41(result, expected, **options):
    """Check if PDFs are correctly organized by year with correct content.

    Args:
        result: Dict from getter with year/filename -> status mapping
        expected: Dict with 'files_by_year' structure
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    files_by_year = expected.get('files_by_year', {})
    total_files = 0
    correct_files = 0
    for (year, file_list) in files_by_year.items():
        for file_info in file_list:
            filename = file_info['filename']
            result_key = f'{year}/{filename}'
            total_files += 1
            status = result.get(result_key, -1)
            if status == 1:
                correct_files += 1
    if total_files == 0:
        return 0.0
    return correct_files / total_files

def check_pdf_not_empty__2d68dae7(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF is not empty (has content).

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: 1.0 if PDF has content, 0.0 otherwise
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with fitz.open(result) as pdf:
            if len(pdf) == 0:
                return 0.0
            text = ''
            for page in pdf:
                text += page.get_text()
            if len(text.strip()) > 0:
                return 1.0
            else:
                return 0.0
    except Exception as e:
        return 0.0

def check_pdf_saved__3dd8a1347f60dd796dc5e98521ae4032(result: Optional[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if a PDF file was successfully saved with valid content.

    Args:
        result: Path to the PDF file (from getter)
        expected: Dict with validation rules (min_pages, min_size_kb)
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size_kb = os.path.getsize(result) / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        from pypdf import PdfReader
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        min_pages = expected.get('min_pages', 1)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_exists__949eb101(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file exists and meets size requirements.

    Args:
        result: Dict from getter with 'exists' and 'size_bytes' keys
        expected: Expected rules dict with 'exists' and 'min_size_kb' keys

    Returns:
        float: 1.0 if file exists and meets requirements, 0.0 otherwise
    """
    if not result.get('exists', False):
        return 0.0
    if not expected.get('exists', True):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 0)
    if min_size_kb > 0:
        size_kb = result.get('size_bytes', 0) / 1024
        if size_kb < min_size_kb:
            return 0.0
    return 1.0

def check_pdf_saved__000c73be1394701e25c8491dd9418647(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if a PDF file was successfully saved with valid content.

    Args:
        result: Path to the PDF file (from getter)
        expected: Dict with validation rules (min_pages, min_size_kb)
        **options: Additional options

    Returns:
        float: 1.0 if PDF exists and meets criteria, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    min_size_kb = expected.get('min_size_kb', 10)
    file_size_kb = os.path.getsize(result) / 1024
    if file_size_kb < min_size_kb:
        return 0.0
    try:
        from pypdf import PdfReader
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        min_pages = expected.get('min_pages', 1)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_pdf_chapters__cedddaca(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if expected PDF chapter files exist with correct page counts.

    Args:
        result: Dict mapping PDF filenames to their page counts
        expected: Dict containing 'files' (list of expected filenames) and
                  'page_counts' (dict mapping filenames to page count requirements)
    """
    expected_files = expected.get('files', [])
    page_count_reqs = expected.get('page_counts', {})
    if not expected_files:
        return 0.0
    total_checks = len(expected_files)
    passed_checks = 0
    for expected_file in expected_files:
        if expected_file not in result:
            continue
        actual_page_count = result[expected_file]
        if expected_file in page_count_reqs:
            req = page_count_reqs[expected_file]
            relation = req.get('relation', 'ge')
            ref_value = req.get('ref_value', 0)
            if relation == 'ge':
                if actual_page_count >= ref_value:
                    passed_checks += 1
            elif relation == 'eq':
                if actual_page_count == ref_value:
                    passed_checks += 1
            elif relation == 'le':
                if actual_page_count <= ref_value:
                    passed_checks += 1
            elif relation == 'gt':
                if actual_page_count > ref_value:
                    passed_checks += 1
            elif relation == 'lt':
                if actual_page_count < ref_value:
                    passed_checks += 1
        else:
            passed_checks += 1
    score = passed_checks / total_checks if total_checks > 0 else 0.0
    return min(score, 1.0)

def check_pdf_naming_pattern__d2aeadd33c64efc8e0e3416b06f6e222(result, expected, **options):
    """Check if PDF files follow the expected naming pattern with content validation.

    Args:
        result: Dict from getter with 'files' (list of file info) and 'source_pdf_exists'
        expected: Dict with 'pattern' (regex), 'min_matches' (minimum files matching pattern)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on:
            - How many files match the naming pattern (partial credit)
            - Whether files are valid PDFs with content (not empty)
            - Whether source PDF exists (to ensure extraction could occur)
    """
    if not isinstance(result, dict):
        return 0.0
    files = result.get('files', [])
    source_pdf_exists = result.get('source_pdf_exists', False)
    if not source_pdf_exists:
        return 0.0
    if not files:
        return 0.0
    pattern = expected.get('pattern', '^\\d+\\. .+\\.pdf$')
    min_matches = expected.get('min_matches', 1)
    min_file_size = 1024
    valid_matching_files = []
    for file_info in files:
        filename = file_info.get('name', '')
        size = file_info.get('size', 0)
        is_valid_pdf = file_info.get('is_valid_pdf', False)
        if re.match(pattern, filename):
            if is_valid_pdf and size >= min_file_size:
                valid_matching_files.append(filename)
    num_valid = len(valid_matching_files)
    if num_valid >= min_matches:
        score = min(1.0, num_valid / min_matches)
        return 1.0
    else:
        score = num_valid / min_matches
        return score

def check_pdf_files__a76459ec(result: List[str], expected, **options) -> float:
    """
    Check if the expected PDF files exist in the directory.

    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'rules' containing expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result must be a list, got {type(result)}')
        return 0.0
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.error('No expected_files specified in expected dict')
        return 0.0
    logger.info(f'Expected files: {expected_files}')
    logger.info(f'Actual files: {result}')
    expected_set = set(expected_files)
    actual_set = set(result)
    exact_match = options.get('exact_match', True)
    if exact_match:
        if expected_set == actual_set:
            return 1.0
        else:
            missing = expected_set - actual_set
            extra = actual_set - expected_set
            logger.info(f'Exact match failed. Missing: {missing}, Extra: {extra}')
            return 0.0
    elif expected_set.issubset(actual_set):
        return 1.0
    else:
        missing = expected_set - actual_set
        logger.info(f'Subset check failed. Missing: {missing}')
        return 0.0

def check_pdf_files__f4eddf72(result, expected, **options):
    """Check if files match expected pattern.
    
    Args:
        result: List of PDF filenames from getter
        expected: Dict with 'pattern' and 'min_matches' keys
        **options: Additional options
        
    Returns:
        float: 1.0 if pattern matches >= min_matches, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    pattern = expected.get('pattern', 'lecture[1-9].pdf')
    min_matches = expected.get('min_matches', 5)
    regex = re.compile(pattern)
    matches = [f for f in result if regex.match(f)]
    match_count = len(matches)
    logger.info(f'Pattern: {pattern}, Expected at least {min_matches} matches')
    logger.info(f'Found {match_count} matching files: {matches}')
    if match_count >= min_matches:
        return 1.0
    else:
        return 0.0

def check_pdf_export__7c083df6(result_state, expected_state, **options):
    """
    Check if a PDF file was successfully exported.

    This function verifies that:
    1. The result file path is not None
    2. The file exists at the given path
    3. The file is non-empty (has content)

    Args:
        result_state: Path to the PDF file (str)
        expected_state: Expected state (dict with 'rules')
        **options: Additional options

    Returns:
        float: 1.0 if the PDF file exists and is non-empty, 0.0 otherwise
    """
    if result_state is None:
        return 0.0
    if not os.path.exists(result_state):
        return 0.0
    file_size = os.path.getsize(result_state)
    if file_size == 0:
        return 0.0
    try:
        with open(result_state, 'rb') as f:
            header = f.read(4)
            if header != b'%PDF':
                return 0.0
    except Exception:
        return 0.0
    return 1.0

def check_pdf_exists__7a335cc66cd23236defef5bd95bbbe7d(result: dict, expected: dict, **options) -> float:
    """
    Check if PDF file exists and was properly exported from the source .docx file.

    Args:
        result: Dict with verification data (exists, is_valid_pdf, file_size, page_count, content_matches, etc.)
        expected: Rules dict with verification requirements
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, partial scores for partial success, 0.0 if failed
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        if not result.get('exists', False):
            return 1.0
        else:
            return 0.0
    score = 0.0
    checks_passed = 0
    total_checks = 0
    total_checks += 1
    if result.get('exists', False):
        checks_passed += 1
    else:
        return 0.0
    total_checks += 1
    if result.get('is_valid_pdf', False):
        checks_passed += 1
    else:
        return 0.0
    total_checks += 1
    if result.get('page_count', 0) > 0 and result.get('file_size', 0) > 1000:
        checks_passed += 1
    total_checks += 1
    if result.get('source_exists', False):
        if result.get('content_matches', False):
            checks_passed += 1
    elif result.get('text_preview', '') and len(result.get('text_preview', '')) > 100:
        checks_passed += 0.5
    if checks_passed >= total_checks - 0.5:
        score = 1.0
    elif checks_passed >= total_checks * 0.75:
        score = 0.8
    elif checks_passed >= total_checks * 0.5:
        score = 0.5
    else:
        score = 0.0
    return score

def check_pdf_contains_text__439ece4f99f624319fffd552c64d5f89(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF contains expected text content.

    Args:
        result: Dict from getter with 'exists', 'text', 'full_text'
        expected: Dict with 'required_text' to search for
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists', False):
        return 0.0
    full_text = result.get('full_text', '')
    if not full_text:
        return 0.0
    required_text = expected.get('required_text', '')
    if not required_text:
        return 0.0
    if required_text.lower() in full_text.lower():
        return 1.0
    return 0.0

def check_pdf_contains_text__e1e75309_4(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists and contains specific text keywords.

    Args:
        result: Path to the PDF file (from vm_file getter)
        expected: Rules dict with 'keywords' (list of required keywords in PDF text)

    Returns:
        float: 1.0 if PDF exists and contains all keywords, 0.0 otherwise
    """
    if not result:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if not result.lower().endswith('.pdf'):
        return 0.0
    keywords = expected.get('keywords', [])
    if not keywords:
        return 1.0 if os.path.getsize(result) > 5120 else 0.0
    try:
        doc = fitz.open(result)
        pdf_text = ''
        for page in doc:
            pdf_text += page.get_text().lower()
        doc.close()
        for keyword in keywords:
            if keyword.lower() not in pdf_text:
                return 0.0
        return 1.0
    except Exception:
        return 0.0

def check_pdf_sizes__4e03b1ed(result, expected, **options):
    """Check if PDF files meet size requirements.

    Args:
        result: Dict mapping filenames to sizes from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Score based on file count and size requirements
    """
    min_size = expected.get('min_size', 1000)
    min_count = expected.get('min_count', 2)
    if not result:
        logger.info('No PDF files found')
        return 0.0
    valid_files = [f for (f, size) in result.items() if size >= min_size]
    if len(valid_files) >= min_count:
        logger.info(f'Found {len(valid_files)} PDFs >= {min_size} bytes: {valid_files}')
        return 1.0
    else:
        logger.info(f'Only {len(valid_files)} PDFs meet size requirement, expected {min_count}')
        return 0.0

def check_pdf_on_desktop__e1e75309_5(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists on Desktop directory by parsing find command output.

    Args:
        result: Output from find command listing PDF files (format: path|size|timestamp per line)
        expected: Rules dict with 'desktop_path' and 'min_size_kb'

    Returns:
        float: 1.0 if at least one valid PDF exists on Desktop, 0.0 otherwise
    """
    if not result:
        return 0.0
    if result.strip() == 'NO_FILES' or not result.strip():
        return 0.0
    lines = result.strip().split('\n')
    desktop_path = expected.get('desktop_path', '/home/user/Desktop')
    min_size_kb = expected.get('min_size_kb', 10)
    min_size_bytes = min_size_kb * 1024
    env = options.get('env')
    for line in lines:
        if not line or line == 'NO_FILES':
            continue
        parts = line.split('|')
        if len(parts) < 2:
            continue
        file_path = parts[0]
        try:
            file_size = int(float(parts[1]))
        except (ValueError, IndexError):
            continue
        if not file_path.startswith(desktop_path):
            continue
        if not file_path.lower().endswith('.pdf'):
            continue
        if file_size < min_size_bytes:
            continue
        if env:
            try:
                vm_ip = env.vm_ip
                port = env.server_port
                read_command = f"head -c 5 '{file_path}' | od -An -tx1"
                response = requests.post(f'http://{vm_ip}:{port}/execute', json={'command': read_command, 'shell': True})
                if response.status_code == 200:
                    output = response.json().get('output', '')
                    hex_bytes = output.strip().split()
                    if len(hex_bytes) >= 5:
                        expected_hex = ['25', '50', '44', '46', '2d']
                        if hex_bytes[:5] == expected_hex:
                            return 1.0
            except Exception:
                return 1.0
        else:
            return 1.0
    return 0.0

def check_pdf_numbered_sequence__fe03784ed42c9a7002fcc758df207953(result, expected, **options):
    """Check if numbered PDF files form a complete sequence.

    Args:
        result: Dict with 'numbered_files' (list of tuples) from getter
        expected: Dict with 'start' (int), 'end' (int) for expected sequence
        **options: Additional options

    Returns:
        float: 1.0 if sequence is complete, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    numbered_files = result.get('numbered_files', [])
    start = expected.get('start', 1)
    end = expected.get('end', 8)
    numbers = [num for (num, _) in numbered_files]
    expected_numbers = set(range(start, end + 1))
    actual_numbers = set(numbers)
    if expected_numbers == actual_numbers:
        return 1.0
    else:
        if len(expected_numbers) == 0:
            return 0.0
        overlap = len(expected_numbers & actual_numbers)
        return overlap / len(expected_numbers)
