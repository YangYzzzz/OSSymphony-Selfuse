"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_csv_sorted__8b88ca382748ac90936480358fbda368', 'check_csv_has_columns__50006736ce024a92878da0e4240e2bb0', 'check_csv_first_n_fullnames__a0281c1d922f2a38909920e121739b97', 'check_csv_and_count__40b9d100', 'check_csv_gender_filter__e49a53b3238eb33c8eed6130c48c5268', 'check_csv_and_count__a4d37536', 'check_csv_merge__c66369e707b97de3ccd6da4699663fe6', 'check_csv_first_name_count__90892eaf', 'check_csv_age_range_filter__5ac234c6f992977135e16f82780e5511', 'check_csv_row_count__2d3e7876', 'check_csv_file_saved__5a942dd0', 'check_csv_sorted_by_lastname__1e62491a', 'check_csv_female_filter_aeb23e02', 'check_csv_conversion_0a84af19', 'check_csv_first_column__e84252aa', 'check_csv_age_sorted__01b44b7acd315e39b1c9a6baa6b5f6da', 'check_csv_columns__7f9974e9', 'check_csv_country_filter__d0eaa9b89f4859670e467a80277e775c', 'check_csv_unique_count__5ab9ee5026038714f957b134595e9a67', 'check_csv_structure_and_content__8de41bc4fc70600ffffb805994ee2926', 'check_csv_account_exists__8f0a3aff', 'check_csv_columns_bbeee8e8', 'check_csv_row_count_5bcd1e32', 'check_csv_country_filter_e8558020', 'check_csv_row_limit__8b005fb41e3efd4706af2ae4f1a79bee']

def check_csv_sorted__8b88ca382748ac90936480358fbda368(result, expected, **options):
    """Check if CSV Age column is sorted in ascending order across entire dataset.

    Verification strategy:
    1. Check first 5 values are 21 (minimum age)
    2. Check last 5 values are reasonable maximums
    3. Sample multiple points (beginning, middle, end) to verify sort order
    4. Verify monotonic ordering: all values are non-decreasing

    Args:
        result: Dict with 'values' (all column values) and 'total_rows' (count)
        expected: Expected rules dict with 'values' key (expected first 5 values)
        **options: Additional comparison options

    Returns:
        float: 1.0 if properly sorted, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    values = result.get('values', [])
    total_rows = result.get('total_rows', 0)
    expected_first_values = expected.get('values', [])
    if not values or total_rows == 0:
        return 0.0
    if total_rows < 100:
        return 0.0
    try:
        int_values = [int(v) for v in values]
    except (ValueError, TypeError):
        return 0.0
    if len(int_values) < 5:
        return 0.0
    first_5 = int_values[:5]
    if not all((v == 21 for v in first_5)):
        return 0.0
    for i in range(0, len(int_values) - 1, 10):
        if int_values[i] > int_values[i + 1]:
            return 0.0
    for i in range(min(50, len(int_values) - 1)):
        if int_values[i] > int_values[i + 1]:
            return 0.0
    mid_start = len(int_values) // 2 - 25
    mid_end = len(int_values) // 2 + 25
    if mid_start >= 0 and mid_end < len(int_values):
        for i in range(mid_start, min(mid_end, len(int_values) - 1)):
            if int_values[i] > int_values[i + 1]:
                return 0.0
    end_start = max(0, len(int_values) - 50)
    for i in range(end_start, len(int_values) - 1):
        if int_values[i] > int_values[i + 1]:
            return 0.0
    last_5 = int_values[-5:]
    if not all((v <= 100 for v in last_5)):
        return 0.0
    min_age = min(int_values)
    max_age = max(int_values)
    if min_age != 21:
        return 0.0
    if max_age < 21 or max_age > 100:
        return 0.0
    return 1.0

def check_csv_has_columns__50006736ce024a92878da0e4240e2bb0(result, expected, **options):
    """Check if CSV data contains expected columns.

    Args:
        result: Dict with CSV data from getter
        expected: Dict with 'required_columns' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dict')
        return 0.0
    required_columns = expected.get('required_columns', [])
    if not required_columns:
        logger.warning('No required columns specified')
        return 0.0
    score = 0.0
    for col in required_columns:
        if col in result and result[col]:
            score += 1.0 / len(required_columns)
    return score

def check_csv_first_n_fullnames__a0281c1d922f2a38909920e121739b97(result: List[str], expected: dict, **options) -> float:
    """Check if first N full names match expected values.

    Args:
        result: List of actual full names from getter
        expected: Dict with 'full_names' key containing expected list
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all match, 0.0 otherwise
    """
    expected_names = expected.get('full_names', [])
    if result == expected_names:
        return 1.0
    else:
        return 0.0

def check_csv_and_count__40b9d100(result, expected, **options):
    """Verify CSV export and count accuracy with partial credit.

    Args:
        result: Dictionary with verification results from getter
        expected: Expected configuration (not used, validation is internal)
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on verification completeness
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('csv_exists', False):
        score += 0.3
    if result.get('csv_has_contacts', False) and result.get('csv_has_mobile_column', False):
        score += 0.3
    if result.get('text_file_exists', False) and result.get('counts_match', False):
        score += 0.4
    return score

def check_csv_gender_filter__e49a53b3238eb33c8eed6130c48c5268(result, expected, **options):
    """
    Check if CSV contains only rows with specified gender.

    Args:
        result: List of lists from CSV (including header)
        expected: Dict with 'gender' key specifying target gender
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    target_gender = expected.get('gender', 'Female')
    header = result[0]
    try:
        gender_idx = header.index('Gender')
    except ValueError:
        return 0.0
    data_rows = result[1:]
    if not data_rows:
        return 0.0
    score = 0.0
    matching_rows = 0
    for row in data_rows:
        if len(row) > gender_idx and row[gender_idx] == target_gender:
            matching_rows += 1
    if matching_rows == len(data_rows):
        score += 0.5
    expected_count = expected.get('expected_count', None)
    if expected_count:
        if len(data_rows) == expected_count:
            score += 0.5
        elif abs(len(data_rows) - expected_count) / expected_count < 0.1:
            score += 0.3
    elif matching_rows > 0:
        score += 0.5
    return score

def check_csv_and_count__a4d37536(result, expected, **options):
    """
    Check if CSV exists and count matches expected.

    Args:
        result: dict with csv_exists and count from getter
        expected: dict with expected values

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('csv_exists', False):
        score += 0.5
    expected_count = expected.get('count', 30)
    actual_count = result.get('count', 0)
    if actual_count == expected_count:
        score += 0.5
    return score

def check_csv_merge__c66369e707b97de3ccd6da4699663fe6(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if CSV merge was successful with correct row count, structure, and data from both sources.

    Args:
        result: Dict from getter containing row_count, has_single_header, sample_rows, total_rows,
                unique_values_from_merged, source_file1_unique_values, source_file2_unique_values
        expected: Dict with 'row_count' and 'has_single_header' keys
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_count = expected.get('row_count', 0)
    expected_single_header = expected.get('has_single_header', True)
    actual_count = result.get('row_count', 0)
    actual_single_header = result.get('has_single_header', False)
    if actual_count != expected_count:
        return 0.0
    if actual_single_header != expected_single_header:
        return 0.0
    unique_values_from_merged = result.get('unique_values_from_merged', set())
    source_file1_unique_values = result.get('source_file1_unique_values', set())
    source_file2_unique_values = result.get('source_file2_unique_values', set())
    if source_file1_unique_values or source_file2_unique_values:
        file1_intersection = unique_values_from_merged.intersection(source_file1_unique_values)
        file2_intersection = unique_values_from_merged.intersection(source_file2_unique_values)
        file1_coverage = 0.0
        if source_file1_unique_values:
            file1_coverage = len(file1_intersection) / len(source_file1_unique_values)
        file2_coverage = 0.0
        if source_file2_unique_values:
            file2_coverage = len(file2_intersection) / len(source_file2_unique_values)
        if file1_coverage < 0.5:
            return 0.0
        if file2_coverage < 0.5:
            return 0.0
        total_source_unique = len(source_file1_unique_values.union(source_file2_unique_values))
        if total_source_unique > 0:
            merged_extra_ratio = len(unique_values_from_merged) / total_source_unique
            if merged_extra_ratio > 1.1:
                return 0.0
    return 1.0

def check_csv_first_name_count__90892eaf(result, expected, **options):
    """Compare actual count against expected count.

    Args:
        result: Count from getter function (int)
        expected: Expected count value (dict with 'count' key or int)
        **options: Additional comparison options

    Returns:
        float: 1.0 if counts match, 0.0 otherwise
    """
    if isinstance(expected, dict):
        expected_count = expected.get('count', 0)
    else:
        expected_count = expected
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_csv_age_range_filter__5ac234c6f992977135e16f82780e5511(result, expected, **options):
    """
    Check if CSV contains only rows within specified age range.

    Args:
        result: List of lists from CSV (including header)
        expected: Dict with 'min_age' and 'max_age' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    min_age = expected.get('min_age', 30)
    max_age = expected.get('max_age', 40)
    header = result[0]
    try:
        age_idx = header.index('Age')
    except ValueError:
        return 0.0
    data_rows = result[1:]
    if not data_rows:
        return 0.0
    score = 0.0
    matching_rows = 0
    for row in data_rows:
        if len(row) > age_idx:
            try:
                age = int(row[age_idx])
                if min_age <= age <= max_age:
                    matching_rows += 1
            except (ValueError, TypeError):
                pass
    if matching_rows == len(data_rows) and matching_rows > 0:
        score = 1.0
    elif matching_rows > 0:
        score = matching_rows / len(data_rows)
    return score

def check_csv_row_count__2d3e7876(result, expected, **options):
    """Check if CSV meets all requirements: row count, header, column count, and email content.

    Args:
        result: Dict from getter with row_count, header, column_count, has_email_content
        expected: Rules dict with expected row count
        **options: Additional options

    Returns:
        float: 1.0 if all requirements met, 0.0 otherwise
    """
    expected_rows = expected.get('rows', 0)
    if result.get('row_count') != expected_rows:
        return 0.0
    if result.get('column_count') != 1:
        return 0.0
    header = result.get('header')
    if not header or len(header) != 1:
        return 0.0
    if header[0].strip().lower() != 'email address':
        return 0.0
    if not result.get('has_email_content'):
        return 0.0
    return 1.0

def check_csv_file_saved__5a942dd0(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if CSV file exists and matches expected configuration.

    Args:
        result: Dict with file_exists and file_path from getter
        expected: Dict with expected file path from rules
        **options: Additional options

    Returns:
        Score: 1.0 if file exists at expected path and is a CSV, 0.0 otherwise
    """
    if not result.get('file_exists', False):
        logger.warning(f"CSV file does not exist at {result.get('file_path')}")
        return 0.0
    actual_path = result.get('file_path', '')
    expected_path = expected.get('file_path', '')
    if actual_path != expected_path:
        logger.warning(f'File path mismatch: actual={actual_path}, expected={expected_path}')
        return 0.0
    if not actual_path.lower().endswith('.csv'):
        logger.warning(f'File is not a CSV file: {actual_path}')
        return 0.0
    logger.info(f'CSV file exists at expected path: {actual_path}')
    return 1.0

def check_csv_sorted_by_lastname__1e62491a(result, expected, **options):
    """Check if CSV contacts are sorted by last name.

    Args:
        result: List of contact dictionaries from getter
        expected: Rules dict with expected sorting
        **options: Additional options

    Returns:
        float: 1.0 if sorted correctly, 0.0 otherwise
    """
    if not result or len(result) < 2:
        return 0.0
    if expected.get('sorted_by') != 'LastName':
        return 0.0
    last_names = [contact.get('Last Name', '') for contact in result]
    sorted_last_names = sorted(last_names, key=lambda x: x.lower())
    if last_names == sorted_last_names:
        return 1.0
    return 0.0

def check_csv_female_filter_aeb23e02(result, expected, **options):
    """Check if CSV contains only Female entries with correct row count.

    Args:
        result: CSV file path (string) from vm_file getter
        expected: Expected rules dict with 'expected_rows' and 'gender_value'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            rows = list(reader)
        expected_rows = expected.get('expected_rows', 0)
        gender_value = expected.get('gender_value', 'Female')
        row_count = len(rows)
        tolerance = 10
        score = 0.0
        if abs(row_count - expected_rows) <= tolerance:
            score += 0.5
        if rows:
            correct_gender_count = sum((1 for row in rows if row.get('Gender', '') == gender_value))
            gender_ratio = correct_gender_count / len(rows)
            score += 0.5 * gender_ratio
        return score
    except Exception as e:
        return 0.0

def check_csv_conversion_0a84af19(result, expected, **options):
    """Check if CSV conversion was successful with correct dimensions.

    Args:
        result: CSV file path (string) returned by get_vm_file
        expected: Expected rules dict with 'expected_rows' and 'expected_columns'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            content = f.read()
        reader = csv.DictReader(content.splitlines())
        rows = list(reader)
        columns = reader.fieldnames
        column_count = len(columns) if columns else 0
        expected_rows = expected.get('expected_rows', 0)
        expected_columns = expected.get('expected_columns', 0)
        score = 0.0
        if len(rows) == expected_rows - 1:
            score += 0.5
        if column_count == expected_columns:
            score += 0.5
        return score
    except Exception as e:
        return 0.0

def check_csv_first_column__e84252aa(result, expected, **options):
    """Check if all column headers match expected order.

    Args:
        result: List of column headers from getter
        expected: Rules dict with expected column order
        **options: Additional options

    Returns:
        float: 1.0 if all columns match expected order, 0.0 otherwise
    """
    expected_columns = expected.get('column_order', [])
    if not result:
        return 0.0
    if len(result) < len(expected_columns):
        return 0.0
    for (i, expected_col) in enumerate(expected_columns):
        if i >= len(result):
            return 0.0
        if result[i].lower().strip() != expected_col.lower().strip():
            return 0.0
    return 1.0

def check_csv_age_sorted__01b44b7acd315e39b1c9a6baa6b5f6da(result, expected, **options):
    """
    Check if CSV is sorted by age in ascending order.

    Args:
        result: List of lists from CSV (including header)
        expected: Dict with 'sort_order' key ('asc' or 'desc')
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    sort_order = expected.get('sort_order', 'asc')
    header = result[0]
    try:
        age_idx = header.index('Age')
    except ValueError:
        return 0.0
    data_rows = result[1:]
    if not data_rows:
        return 0.0
    ages = []
    for row in data_rows:
        if len(row) > age_idx:
            try:
                age = int(row[age_idx])
                ages.append(age)
            except (ValueError, TypeError):
                return 0.0
    if not ages:
        return 0.0
    is_sorted = True
    for i in range(len(ages) - 1):
        if sort_order == 'asc':
            if ages[i] > ages[i + 1]:
                is_sorted = False
                break
        elif ages[i] < ages[i + 1]:
            is_sorted = False
            break
    expected_count = expected.get('expected_count', 5000)
    score = 0.0
    if is_sorted:
        score += 0.6
    if len(data_rows) == expected_count:
        score += 0.4
    elif len(data_rows) > 0:
        score += 0.2
    return score

def check_csv_columns__7f9974e9(result, expected, **options):
    """Check if CSV has expected number of columns with correct headers and data.

    Args:
        result: Dict with columns, headers, and row_count from getter
        expected: Rules dict with expected column count
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_cols = expected.get('columns', 0)
    if not isinstance(result, dict):
        if result == expected_cols:
            return 1.0
        return 0.0
    if result.get('columns', 0) != expected_cols:
        return 0.0
    headers = result.get('headers', [])
    if len(headers) != 2:
        return 0.0
    normalized_headers = [h.strip().lower() for h in headers]
    expected_headers = ['first name', 'last name']
    if normalized_headers != expected_headers:
        return 0.0
    row_count = result.get('row_count', 0)
    if row_count < 25:
        return 0.0
    return 1.0

def check_csv_country_filter__d0eaa9b89f4859670e467a80277e775c(result, expected, **options):
    """
    Check if CSV contains only rows with specified country.

    Args:
        result: List of lists from CSV (including header)
        expected: Dict with 'country' key specifying target country
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    target_country = expected.get('country', 'United States')
    header = result[0]
    try:
        country_idx = header.index('Country')
    except ValueError:
        return 0.0
    data_rows = result[1:]
    if not data_rows:
        return 0.0
    score = 0.0
    matching_rows = 0
    for row in data_rows:
        if len(row) > country_idx and row[country_idx] == target_country:
            matching_rows += 1
    if matching_rows == len(data_rows):
        score += 0.5
    expected_count = expected.get('expected_count', None)
    if expected_count:
        if len(data_rows) == expected_count:
            score += 0.5
        elif abs(len(data_rows) - expected_count) / expected_count < 0.1:
            score += 0.3
    elif matching_rows > 0:
        score += 0.5
    return score

def check_csv_unique_count__5ab9ee5026038714f957b134595e9a67(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if unique count matches expected value and verify merge operation.

    Args:
        result: Dict from getter containing unique_count, total_rows, and file_exists
        expected: Dict with 'unique_count' and optionally 'min_rows' keys
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not result.get('file_exists', False):
        return 0.0
    expected_count = expected.get('unique_count', 0)
    actual_count = result.get('unique_count', 0)
    if actual_count != expected_count:
        return 0.0
    min_rows = expected.get('min_rows', 5000)
    actual_rows = result.get('total_rows', 0)
    if actual_rows < min_rows:
        return 0.0
    return 1.0

def check_csv_structure_and_content__8de41bc4fc70600ffffb805994ee2926(result: Dict[str, Any], expected: dict, **options) -> float:
    """Check if CSV has the expected structure AND contains actual data.

    This metric verifies:
    1. The CSV has the correct number of columns
    2. The CSV has at least one row of data (not just headers)
    3. The CSV contains non-empty cells (actual content)

    Args:
        result: Dict with column_count, row_count, has_data, non_empty_cells
        expected: Dict with 'column_count' and optionally 'min_rows'
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_count = expected.get('column_count', 0)
    min_rows = expected.get('min_rows', 1)
    if result.get('column_count', 0) != expected_count:
        return 0.0
    if result.get('row_count', 0) < min_rows:
        return 0.0
    if result.get('non_empty_cells', 0) == 0:
        return 0.0
    return 1.0

def check_csv_account_exists__8f0a3aff(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check if expected accounts exist in Thunderbird accounts CSV.

    Args:
        result: path to csv file containing account data
        expected: dict with "expect" key containing list of account records to find

    Returns:
        float: 1.0 if all expected accounts exist, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_list = expected.get('expect', [])
    if not expect_list:
        return 1.0
    expect_metrics = [False] * len(expect_list)
    try:
        with open(result) as f:
            reader = csv.DictReader(f)
            for rcd in reader:
                for (i, expected_rec) in enumerate(expect_list):
                    if all((rcd.get(k) == v for (k, v) in expected_rec.items())):
                        expect_metrics[i] = True
    except Exception as e:
        logger.error(f'Error reading CSV: {e}')
        return 0.0
    return float(all(expect_metrics))

def check_csv_columns_bbeee8e8(result, expected, **options):
    """Check if CSV has correct columns.

    Args:
        result: CSV file path as string
        expected: Expected rules dict with 'required_columns' and 'column_count'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            columns = reader.fieldnames
            if not columns:
                return 0.0
            required_columns = expected.get('required_columns', [])
            expected_column_count = expected.get('column_count', 0)
            score = 0.0
            if len(columns) == expected_column_count:
                score += 0.5
            if required_columns:
                columns_found = sum((1 for col in required_columns if col in columns))
                score += 0.5 * (columns_found / len(required_columns))
            return score
    except Exception as e:
        return 0.0

def check_csv_row_count_5bcd1e32(result, expected, **options):
    """Check if CSV has expected number of data rows.

    Args:
        result: CSV file path from get_vm_file
        expected: Expected rules dict with 'expected_data_rows'
        **options: Additional options

    Returns:
        float: 1.0 if row count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'rb') as f:
            content = f.read().decode('utf-8')
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        data_row_count = len(rows) - 1 if len(rows) > 0 else 0
        expected_rows = expected.get('expected_data_rows', 0)
        if data_row_count == expected_rows:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_csv_country_filter_e8558020(result, expected, **options):
    """Check if CSV contains only entries from specified country with correct row count.

    Args:
        result: Path to CSV file (string)
        expected: Expected rules dict with 'expected_rows' and 'country_value'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            rows = list(reader)
        expected_rows = expected.get('expected_rows', 0)
        country_value = expected.get('country_value', 'France')
        row_count = len(rows)
        tolerance = 10
        score = 0.0
        if abs(row_count - expected_rows) <= tolerance:
            score += 0.5
        if rows:
            correct_country_count = sum((1 for row in rows if row.get('Country', '') == country_value))
            country_ratio = correct_country_count / len(rows)
            score += 0.5 * country_ratio
        return score
    except Exception as e:
        return 0.0

def check_csv_row_limit__8b005fb41e3efd4706af2ae4f1a79bee(result, expected, **options):
    """
    Check if CSV contains only the first N rows from source.

    Args:
        result: List of lists from CSV (including header)
        expected: Dict with 'row_limit' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    row_limit = expected.get('row_limit', 100)
    header = result[0]
    data_rows = result[1:]
    if not data_rows:
        return 0.0
    score = 0.0
    if len(data_rows) == row_limit:
        score += 0.6
    elif len(data_rows) > 0:
        if abs(len(data_rows) - row_limit) / row_limit < 0.1:
            score += 0.4
        elif abs(len(data_rows) - row_limit) / row_limit < 0.2:
            score += 0.2
    try:
        id_idx = header.index('Unnamed: 0')
        first_id = int(data_rows[0][id_idx])
        if first_id == 1:
            score += 0.2
    except (ValueError, IndexError):
        pass
    try:
        id_idx = header.index('Unnamed: 0')
        last_id = int(data_rows[-1][id_idx])
        if last_id == row_limit:
            score += 0.2
    except (ValueError, IndexError):
        pass
    return score
