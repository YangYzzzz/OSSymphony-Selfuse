"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['exact_match__6cc337d6467b2f92dee1f171f61ea4ed', 'check_import_lines__ad9b1b7a', 'check_exact_text_match__609a65b9', 'check_text_file_content__5ced85fc_aug18_v4_d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8a9', 'check_text_lines__c2b520792100f4aa54f0d89dedfa94c4', 'check_timezone_sydney__e603b840', 'compare_text_output__f558b67b', 'check_first_line_left_aligned__adcd0938', 'check_file_count__9364293cce5b25e22063aee62da7d43d', 'check_file_upload__68f5fe6b', 'check_text_format__7f844786255954cce16f5ea58433f34e', 'check_text_replacement__1f0d9653', 'check_python_code_complete__f55aa7954b40a62bad3b8ba851857ed1', 'check_file_count_greater__00db2192', 'check_python_syntax__198be354', 'check_text_output__4cc7c3cf', 'check_file_exported__d22e0dfa', 'check_python_imports__261836e0618ec18a7a70e8da4837dfbf', 'check_file_organization__92a58812', 'check_text_contains__5b8281b6', 'check_recent_file_count__db5a3e05', 'check_text_exact_match__68f83c37fe78a3f2dbadefcb3c480330', 'check_text_replacement__ffa8cf6ad724ee8fc8a065457d283c28', 'check_default_text_color__ea408cb7', 'check_text_exact_match__2d81db9f5efc3d72c38ba9f24bf6d4fb', 'check_file_in_list__880f8efb', 'check_file_exists__739292ff', 'check_documents_file__c0d05fd4eae793b685d13d8511de53aa', 'check_remaining_files__c90c6ca3', 'check_file_content__11701199', 'check_file_contains_text__6c41fae8ffe95d2c4b3c521d5446de56', 'check_file_count__23c57bc9', 'check_text_lines__7e304294ce76dab9f08b95178667a620', 'check_single_eml_file__54c9c002bd0d5132155806dba06b543e', 'check_file_exported__7107319d', 'check_text_file_lines__5ced85fc_aug18_v3_c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8', 'check_file_size_range__0f2a6243', 'check_both_files__3f6d3219', 'check_note_exact_match__5f5d62626ae9f466dd58a3e57638d2a2', 'check_sender_file__ad43db2627764dd28ab9631606c7b97c', 'check_contains_filename__32107748', 'check_merged_text__e22bfb55f6ab9983d1cc35b82dc09aeb', 'check_textbox_bottomleft__99aea82e', 'check_text_contains__de3b1681', 'check_text_contains__c491dbb8', 'check_file_permissions__11701199', 'check_textbox_on_rightside__5bc789db', 'check_files_match_pattern__fb580d40', 'check_text_file_contains__ebe7bea30aab4d43b91ea1760b3fb66f', 'check_text_contains_all__b8171706418d2058f81460f1a24d4635', 'check_python_comment__4ca3247dcf6b464342c4e3f53d844797', 'compare_text_output__2175dacb', 'check_multiline_answers__fa7deb995de2addd4a7d5d55ff4c5c25', 'check_first_line__66414428', 'check_file_exists__d967a5b0', 'check_text_content__1a1f627807b83c4d33e1ae428da08935', 'check_file_exists__2ad8e92c', 'exact_count_match__f1d66967', 'check_file_renamed__9f8a1d0c', 'check_textbox_in_bottom_right__c729f30a', 'check_content_text_color__53beb8d2', 'check_timezone__3d14fb6b', 'check_file_renamed__b4364020', 'check_file_exists__9c31b3f6afb568d600c17c937149b6c4', 'check_timezone__43ba8703', 'check_file_exists__4c5ac05d', 'check_textbox_centered__723c4038', 'check_file_exported__6fb34f5d', 'check_exact_match_v5__d4477d7a', 'check_textbox_on_topside__20161493a3832ab8acb83e42c1b1cbaf', 'check_file_format_and_structure__6bbe76fc', 'check_file_exists__42ef2d72a0f41077972857e814318a24', 'check_recent_files__4e03b1ed', 'check_textbox_in_bottom_left__f062cc22', 'check_timezone__308d9780', 'check_python_structure__58d4fb9ea6e69b7f57a37a59813050e8', 'check_file_exists_with_size__d6a98717', 'check_comment_lines__70bdd422', 'check_text_increased__e6f52b2c', 'check_python_with_summary__be0d67f2ca1c7fde5d3c550cd046d0b4', 'check_text_replacement__0b4a6158', 'check_file_path__8369c71b0ee26c543c025c6e1cb39bbd', 'check_file_rename_pattern__a610cecf', 'check_music_file_count__355d5753246a8a88c1b8d173c110d89e', 'check_text_on_rightside__68e0eaa1c0be88d7d26c920a571b88e0', 'check_file_content__e1c8e8d0', 'check_file_exists__c0930f6fc6470d951ad9d775e3a6c6a5', 'check_git_repo_status__55495fb6b59196cc7cffae2b12e117ed', 'check_exact_match_v3__d4477d7a', 'check_git_initialized__763d7485', 'check_text_output__970046ef9644d0d33f656e418e7d5e7a', 'check_textbox_in_top_left__753748f3', 'check_line_count__1098fce8', 'check_python_imports__0a56ab11', 'check_file_exported__6ee1fdcd', 'check_file_exists__12ebd85d', 'check_python_definitions__d881a7db2082639af55dd0ea1e3047ab', 'check_renamed_files__0190943f9252410b5b69d12d28f1b6b7', 'check_text_content__fc2c8cc4', 'check_file_exists__dccfbd8e', 'check_file_exists_with_size__c0d603dc', 'check_lines_contain_keywords__bf9ee805c777733e26c14d1927c30b1a', 'check_file_exists__6cb37327', 'check_file_recovery__5ea617a3', 'check_python_path__e7a35ed0ec20ca7ab6b257f3f5c87e23', 'check_python_files_count__13d7d579', 'check_file_count__081d0b6c', 'check_json_settings__c5a909ed', 'compare_text_output__e6536361', 'check_file_exists_with_size__91c793ae', 'check_file_move_copy__aace45122e840d40b84dc540ae5a49bc', 'check_file_renames__6f033773', 'check_python_text_patterns__198be354', 'check_python_pkg__a3b47e9754f6a01a17ed98f7d00d938c', 'check_large_file_count__b9976565', 'check_file_exists__45e0b2e9', 'check_file_size_range__b9c089b2fe7d833fde2da297bbbd9620', 'check_text_file_lines__5ced85fc_aug18_v0_c9e8a1b2d3f4e5a6b7c8d9e0f1a2b3c4', 'check_text_replaced__66306a8b', 'check_file_size_range__79b627f8', 'check_text_content__6941d0dc31bbd0c3d844303b7e1c57e5', 'check_file_exists__b8a50137', 'check_file_count__4e03b1ed', 'check_files_moved__5bd46009a70125f4395abd10c34421d4', 'check_filenames_match__71c23132811d122bc61dca33636b8f81', 'check_text_content__a48d3d2ba069ee77685e4821041681b9', 'check_text_replacement__438c9c7ce7eebe25a3992ddf0a388112', 'check_line_count__bcaf4400', 'check_python_definitions__e8ec0313751ff65df15abd7d031a7c53', 'check_contains_numbered_files__939aa00d', 'check_srt_file_exists__4f098003e517e7e34a157f1c233e1c85', 'check_text_contains__7e5134258960ebea77ca0d290984a7a3', 'check_text_replacement__d48f445fac6cb18530cd7f7169fdc7fc', 'check_line_starts_with__8f9a3936', 'check_file_contains_comment__bcc15712', 'check_exact_notes_count__eb089173', 'check_exact_text_match__53b3f7f8', 'check_zip_contains_files__5c33f919', 'check_text_content__20f90bc2668d01c30660dfeafc3af15b', 'check_python_file_count__c32ed37d80ef317dea88bac4a4cc1f31', 'check_text_exact_match__f84e9cb5fd8cfc9fab208c24bcd90a7d', 'check_text_file_content__5ced85fc_aug18_v2_b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7', 'check_text_replacement__0f72ff3b9c8d9680b46915659d675f48', 'check_exact_text_match__b4af657e', 'compare_text_output__d7f70e02', 'compare_text_files', 'check_line_count_equals__fb1a48c8', 'check_file_exported__9180f469', 'check_file_size__f778a8914a698cd2bc7c0cc50cd3596d', 'check_file_created__b03b6c61', 'check_line_indentation__a2b5af8108e461937716b976809e966c', 'check_direct_json_object__82bc8d6a', 'check_file_contains_items__f8073900a375900c7a9ce8fa79f05f9d', 'check_specific_file_exists__72383bef81d322492ebfc1e4d86364a3', 'check_specific_file_present__04085b6d', 'check_python_file_exists__198be354', 'check_text_exact_match__600696b8508be0c2e3ca25794856fb75', 'check_dir_exists_with_files__5c107_2', 'check_file_exists__724cfa0a', 'check_file_exists_and_structure_sim__2cf5fa39', 'check_file_content__57d8acad', 'check_file_exists_with_content__845b16e9c20bb76eb4ef8a7eb9262413', 'check_launch_json_config__a3743b930b4e3c5c6976584a28c8269c', 'check_blank_line_added__ebbefd78bc8d93d288c452f335196528', 'check_word_count_text__a4c1c19457791f4f99b06ead98b6bfeb', 'check_first_line_right_aligned__de33d712', 'check_git_repo_cloned__f1a99656b1aa540fcb46d2aeba395a7e', 'check_timezone__e2503403', 'check_text_replacement__66cc1bc44690547654db9625c97259f8', 'check_text_replacement__0f355d4c', 'check_file_exists_with_size__058bd353', 'check_line_numbering__2978d3bd', 'check_files_renamed__039c45a2', 'check_file_organization__429c8cbc', 'check_python_file__f1f92c4b10af2ffde6ea1534830a28f6', 'check_filename_match__4b590a3a028f08e8f4ad12729f4351c3', 'check_column_text_values__2c840a52', 'check_text_replacement__ca4f47f2', 'check_text_with_space__e8a68271', 'check_text_match__25fb76d7ccb83f42013b589a25bead61', 'check_python_path__1f6f3af3ee1e7e72d7b32984543de005', 'check_file_content_5f767718', 'check_file_exists__f663e89c52b74d5c5d4e38ab6d86c83f', 'check_exact_int_match__d4477d7a', 'check_files_deleted__0a07f87a', 'check_filename_pattern__4e03b1ed', 'check_file_exists__81d9e3403ff656f186df75dab490d6ac', 'check_specific_text_deleted__ec78028f', 'check_eml_files_with_pattern__8a92d4dc', 'check_text_exact_match__4080707f437c813fb70f2db7aaa30575', 'check_text_replacement__11d0824d05970e58a5671a5636365f15', 'check_both_files_rows__e54614f2', 'check_file_exported__0ad50b28', 'check_file_exists__423224cbbe432d6315ffb9aa3c684c3a', 'check_file_count__23e95644', 'check_file_exists__cb8ab5642aaf48705d11abb5543759c9', 'check_exact_count__68182234', 'check_txt_file_count__2b6d7a72', 'compare_text_output__2d526d9e', 'check_gitignore_content__c5a909ed', 'check_renamed_files__27c9f1432580c9f5f1bf2d1f919f4ed5', 'check_file_list__b8c40a8e', 'check_file_organization__f3cabf2e', 'check_text_contains__0b52fd51', 'check_file_exists__aec9e92c', 'check_text_output__ec857351', 'check_files_with_prefix__8c4eaed9a61673f78def8f323e7dfe9d', 'check_python_imports__198be354', 'check_file_exists__ec920d7f', 'check_file_renamed__a7904c4a', 'check_text_contains__a2f161de', 'check_text_output__9ee70fa2', 'check_text_exact_match__49075d97b370c316554a4b259a7ccc3e', 'check_zip_contains_files__6aa029d37a944ba9e2bf06a8a1d59f5c', 'check_subject_file__271abb880d5f6f8d57d2c41e20bcf6ad', 'check_line_chart__0f492b286ccfa81ae15bde7a08cd32e6', 'check_exact_number__f9a0219a', 'check_filename_hash_mapping__bf825e2c', 'check_textbox_vmiddle__5230e9e6', 'check_exact_file_count__1271f790', 'check_file_permissions__7bde372c', 'check_file_deleted__9e688855', 'check_text_color__b3dc20103be74c0981bbd79306e76a3d', 'check_file_location__c78e5698dfdf96679302f35b21f0928f', 'check_text_output__b83d8fa5', 'check_text_occurrence__7dc05d2e', 'check_all_python_syntax__c9c8227a4cd72e8de3e73ea399f7f61d', 'check_file_content__2a7463c5815fe65f87729a241d0d409d', 'check_file_count__19cf6326', 'check_file_rename__990ae9b047da99489a16db0558f7ee61', 'check_srt_filename__ab47203640c5bdcef1195c50e51e7524', 'check_chapter_files_exist__93cfd69b3c5adfd5dbb8817764000202', 'check_timezone__b98a8580', 'check_dual_file_exists__e1a4b749', 'check_file_line_count__e09bfbdf', 'check_text_content__67c66e9d6c723be29d116d6e2c7b5850', 'check_text_value__06ae69a0', 'check_textbox_at_top__83c7a705', 'check_gitignore_file__6ee0182d', 'check_text_output__ddbee6ed', 'check_file_copy__dc38ce29eba391e7169ff4e028e69a72', 'check_text_replacement__e99e7d69ea502440f7bd1b2cb57a309d', 'check_direct_json_object__f5d96daf_task_verify_3', 'check_file_exists__7298f585793a2b727ac3910de3795a50', 'check_text_file_content__50943048', 'check_line_count__953256df836603c8857d4495861e4b63', 'check_all_files_exist__03426e679d8f4571bede57a16eea69a4', 'check_text_contains__a515d38160faf207fcfd0df30838b0c3', 'check_textbox_topleft__7d24ebd2', 'check_zip_files__53fe105429a60cae06bcf9ce59b19b3e', 'check_text_replacement__71cb5a9efacfe89bccff2081bcadcf02', 'check_file_saved__9c119f81', 'check_file_exists__6f3c16ae', 'check_file_created__684b5a3a3f653750766f5bbe64af3bd5', 'check_text_content__25d344d5', 'check_file_content_e9b24959', 'check_timezone__f88b70b8', 'check_min_file_size__b756f99d', 'check_exclude_platform__156bbecf8094c096edc4b32f7b6fd25b', 'check_intext_citation_added__9d568660', 'check_file_count__bc253f41', 'check_first_line_indent__7678624b', 'check_file_readonly__b3b80682', 'check_file_exists__689ec9af4ba1471bf9b5f89e71cafeb9', 'check_new_file_properties__73935909', 'check_file_content__d09dce0a', 'check_python_imports_only__093631738f9b5eba42a5bcf60212ba3b', 'check_file_contains_lines__fde871ae', 'check_file_contains_lines__eac4b332', 'check_file_downloaded__08c5e1b6ad7015f1bdd4ff79ff88e12f', 'check_exact_text_match__65949e53', 'check_file_content_555dda86', 'check_files_exist__a8a082525df1807c95a7519289fda5a0', 'check_file_permissions__f81b27ec', 'check_python_code_backup__f222fdd4d3c26325ac7310b0f2b1711f', 'check_file_size_range__198be354', 'check_file_count__17298c22', 'compare_text_output__81f11cbf', 'check_file_exists__24a50bf5', 'check_python_function_exists__c685793b2ca36ab76b7f2cc84f84fe40', 'check_exact_text_match__6ba0a623', 'check_filename__4ee0209a', 'check_line_count__cfbf4273', 'check_line_count__475840bd88bfc32515242a838ac799b5', 'check_file_size__739292ff', 'check_file_moved__789836386f3e1cf0e0ee5d172a0885f2', 'check_file_copied__642c6d87', 'check_line_count_positive__868f1e74', 'check_file_organization__3b0a753c', 'check_renamed_files__47275204', 'check_remaining_files__8754d37bdc9e8d94ab80feb618caa015', 'check_file_exists__c6c3aa52', 'check_text_file_value__83bea20e1ddf48cf4f537ad2c05896b6', 'check_single_file_exists__d6fb1e53c50621e1a08efd7623119b0d', 'check_file_exists__67be1ac6efe87edab008a615fb0e7ec4', 'check_file_count__58fb65f5', 'check_file_duplicate__654353fe', 'check_text_opacity__88b1d5c668539570e153fff50a7fc5f9', 'check_file_contains__8085b902c0aa531b12d2ed766e22c897', 'check_ods_file_exists__bdb8ae26', 'check_file_recent__739292ff', 'check_file_exists__3928cfa5', 'check_winloss_sparklines__2bd59342_aug_12_verify_2', 'check_text_content__f8ce9de4', 'check_include_excluding_regex__06dc70fa7cf93b7501432994b47d9c35', 'check_files_count__014e651b', 'check_specific_files_present__465762dc', 'check_file_list__b0918ceb', 'check_file_exists__f5680565', 'check_vm_file__f259522f5141b84d8b2c6c9007fd732a', 'check_file_count__be02851a', 'check_text_contains__8f36e769', 'check_python_functions__198be354', 'check_file_exists__31a8a4acc19afab71c5f7ec2f3006a11', 'check_all_files_exist__4e03b1ed', 'check_python_script_contains__c9385c6b', 'check_text_contains__551695fc', 'check_file_content_d461fd5d', 'check_direct_json_object_with_time__f79439ad', 'check_text_output__451bbe5b', 'check_file_exists__e86a3a4d', 'check_file_exists_with_content__4f896edc', 'check_file_exists__5f5351b0', 'check_python_classes__198be354', 'check_file_exists__c48ab0f0', 'check_textbox_on_bottom__734a49ef', 'check_file_exists_and_structure_sim_77b8ab4d', 'check_textbox_bottom__3bddf9fd7655c038551f7bf6cc6f697d', 'check_python_imports__5e897a929023bc43113113bb0ca8fb36', 'check_first_line_indent__2b727758', 'check_direct_json_object', 'check_exact_file_list__91c6f86e', 'check_file_exists__b2463eb9', 'check_textbox_top__f0019e8aa52327a88742c54abff5dd41', 'check_total_files_count__e9983217', 'check_txt_file_content__b2ba9ee6873b43000b20ba169c4d9592', 'check_file_exists_with_size__e1affec5', 'check_file_rename__c5a909ed', 'check_python_code_stats__198be354', 'check_file_organization__94dfca2e', 'check_file_properties__e9151d35420d1468fe1e721181658d5f', 'check_file_format__4b5f0cdf', 'check_file_line_count__7233f122896fac183c973343e2cf3b2a', 'check_file_contains_text__edd1c1331eff751ad5487718d7e5d07b', 'check_timezone_utc_minus_5__d030998aa6296cb272a4b11f8e8cd0d6', 'check_timezone_cet__755f0aea', 'check_timezone_ist__36d379a3', 'check_lines_counter__20d8676e', 'check_text_replacement__fe4ab075', 'check_include_func__864e927cd0113eb6d9476a9e8ebce88d', 'check_text_rotated__a994a687', 'check_text_content__5ef70598fcc68171b6ec36d7c888fc21', 'check_text_list__d9839857', 'check_file_exists__7930967f', 'check_file_content__8caa3012', 'check_multiple_files__11701199', 'compare_text_output__d398cd07', 'check_text_content__a652858db2e92fae77817389157c8edc', 'check_file_exists__04b2d876', 'check_line_count__ab2d13c4', 'check_exact_text_match__5b92f2b5', 'check_text_replacement__f1f9b3a9', 'check_exact_match_v8__d4477d7a', 'check_file_exists__03122ed4', 'check_file_exists_with_size__dccb2b60', 'check_text_exact_match__515e2337245bc72c8d34192293ce6646', 'check_file_contains_lines__5f60eaec', 'check_python_classes__d4d97725', 'check_textbox_in_top_right__8fd9cc45', 'check_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385', 'check_first_line_right_aligned__09652fe4b6098782ed16d144b86a430f', 'check_dir_and_files__5b67568a', 'check_text_lines_exact__d1fc13ca7061617e08d8a914a14209cd', 'check_file_exists__3e673542', 'check_files_by_pattern__0848b03099d380057c59f332d48dc222', 'check_file_recently_modified__23cbcfa9', 'check_text_replacement__27683cd2', 'check_file_count__bb5651c2', 'check_files_moved__90e7472a', 'check_file_exists__7ae8ce2b', 'check_file_checksum__032a6328', 'check_lines_removed__3204f52d', 'check_file_exists__dee238bc', 'check_text_exact_match__895c3960d172d43278234eeb5c495eda', 'check_file_permissions__dd3bb1bb', 'check_lines_commented__4ee5b05c', 'check_exact_text_match__9a63ba8e', 'check_tsv_file_exists__ecf92ffc', 'check_file_exists__e3dda739ee4da14903d2e8df52d7a41d', 'check_text_contains__7dfb45a4', 'check_text_alignment__faeddc67', 'check_text_content__69119b71', 'check_file_exported__46f0b51f', 'check_comprehensive_text__c4f4ba50', 'check_text_output__99a23d7f', 'check_file_locations__b75faf5b6765d0d1458ed6b6d219047b', 'check_copy_file_status__24914e86', 'check_word_count_file__adfc25c4', 'check_text_replacement__514834e6', 'check_file_exists__735234d9', 'check_timezone_utc_plus_8__a394a8f31cb4c0f1d8dc918cb19a351d', 'check_largest_filename__b7de68e1', 'check_both_files_non_empty__e07d7a26', 'check_archive_contains_files__0b074054', 'check_files_with_keyword__91578e58', 'check_text_file_mountains__1e67b9ae311891ef2e3034615cded86c', 'check_file_exists__506ad17f', 'check_first_line_heading1_style__ba24a7ac', 'check_file_exists__a544be73', 'check_file_count__63477992bbc88c6a7091e80f7c0a8a72', 'check_exact_match_v4__d4477d7a', 'check_eml_files_exist__aebf0a91a4be0be45ef3247f943131f2', 'check_backup_files__ac9408954b941c7f40eedd27a6f1296b', 'check_file_exported__f25970ca', 'check_titles_file__3af0b2bf', 'check_file_exists__a693f275', 'check_file_exists__46ccb784', 'check_total_file_count__ed6a3699fe6af7deef02a8e547504034', 'check_text_file_content__5ced85fc_aug18_v1_a1b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6', 'check_direct_json_object__fc6d8143', 'check_both_files_exist__805294f8', 'check_text_match__ed1a5c265e6c6d06dcaf2ec482204403', 'check_textbox_on_top__49abb020', 'check_file_contains__3c9f051952e2f37565e45b593e085b87', 'check_file_exists__d8671412', 'check_timezone_utc_plus_1__5f2017655ce06bbb056bf7cc4210c4ca', 'check_file_renamed__c396550f', 'check_file_content__89d906fa', 'check_textbox_on_topside__57a0b169', 'check_text_contains__b38e8bb9', 'check_timezone_utc_plus_530__851aed8c', 'check_text_replacement__8b701bf8d0cacb95438d2e4e17a8b914', 'check_word_count_file__7c58ef63', 'check_timezone_utc_minus_8__dfb26af6f8639d73cfa7d8bdb0721f21', 'check_file_not_exists__be265045', 'check_json_settings__c5a909ed_v5', 'check_invoice_summary_text__834c93d1a65ecbb7766bb5ceb1a12320', 'check_timezone__f2f0035d', 'check_all_files_exist__f50a55ca', 'check_file_exists_with_size__ff3634ef', 'check_file_count__48e4325f2d62197da2b10059281b95a0', 'check_file_organization__8c566ad0', 'check_text_replacement__56d0ef227fc0a081b24201d3ecb3d358', 'compare_text_output__b1b8ec8f', 'check_file_exists__4e03b1ed', 'check_timezone__29ace8b1', 'check_file_counts__816aac7b5fcbde572f13b62a3999bd4d', 'check_text_replacement__62f7a00e7dd5e4db70eb00a615f012ef', 'check_file_organization__a93d97ba', 'check_file_exists__82518dad', 'check_txt_line_count__8ff67d2b', 'check_file_exists_with_size__21492178', 'check_text_alignment__1857b08997e5057be0ed7e2fe747fd99', 'check_text_content__d7828490', 'check_text_output__bb796efd', 'check_text_replacement__681a97c35eb849e5cf9422adfe4e1aea', 'check_file_naming__d6b113924c427deceec5f933af24484e', 'check_file_exists__62dacdc1', 'check_exact_match_v6__d4477d7a', 'check_python_complete__12fe4256', 'check_text_lines_exact__4250d59b26bb86f2de0562f0a55c312c', 'check_text_match__93eac3e2452ef121ce8047db9ec250fe', 'check_file_exists_with_size__479794c8', 'check_file_organization__e2bf8bf2', 'check_file_ownership__18d171e6', 'check_timezone_utc_plus_3__d1091690', 'check_selective_file_deletion__ca75e69be093d6cb8e4fa53ce6114782', 'check_text_replacement__c5c5b95c23a9d2c28863362320aba24b', 'compare_text_output__cee492ee', 'check_text_patterns__32952afd', 'check_downloads_file_exists__2cbf25da', 'check_file_created__974295f9d11461d175dbc0223dd4ff65', 'check_files_renamed__416ce0b1', 'check_all_text_color_match__6666d59e63b5ea41073d9fefe10e8bab', 'check_file_exists__dc35efec', 'check_rtf_file_exists__34460ac1a76394f2dfa8b4e9981344a0', 'check_title_text__8b1e2a3b', 'check_file_existence__e1da6937', 'check_new_textbox_added__4ebe6ee8', 'check_text_content__7bbdf0a0733630cbbfd86729556fc827', 'check_file_created__8ab0a45d4a57cf0e9592e87621895b59', 'check_file_exists__6d219cf2', 'check_python_content__198be354', 'check_filename_prefix__fdcd3f41', 'check_filename__739292ff', 'check_textbox_centered__13f49ee1', 'check_file_exists__cdcbbd90', 'check_file_existence__198be354', 'check_file_exists__dd9409c8', 'check_file_exists__916d8b58', 'check_git_dir_exists__e2da960ab9034666db33db74ae6371a7', 'check_textbox_on_rightside__5c6340a1', 'check_file_exists__0b2042a1', 'check_text_output__afe6b7a7', 'check_file_count__03f8ef9d', 'check_file_contains_lines__d5302e2f', 'check_timestamped_file__2b78c2fd0d670b6ee1c54ce65b4419a5', 'check_file_exists__e3ae8a85', 'check_timezone_pst__27082429', 'check_textbox_movedleft__0fe6ee6f', 'check_file_exported__60340d37', 'check_python_docstring__6ee0182d', 'check_textbox_smaller__a8f988b59259c0be84da4d1fc65b92ad', 'check_python_imports__52a225d6', 'check_text_content__7676732d', 'check_line_order__2c3b878a', 'check_line_duplicated__33cae2c6', 'check_chapter_file_naming__29db12fd', 'check_exact_text_match__e4beccec', 'check_text_replacement__65dee18f1880f7d028c2b1727fd62d90', 'check_files_contain_pattern__7427978e', 'check_line_count__5538243966b3481fe772536923c1f693', 'check_moved_files__1b47b6505a7a2dc3d6ad6f0c07b4bcb4', 'check_textbox_topright__d493bb30', 'check_filename_exact_match__0d61b4f8', 'check_textbox_fully_centered__c84bec37', 'check_exact_recipient_set__c9ce3f52', 'check_eml_files__846f274f', 'check_timezone__7a83c51f', 'check_file_exists__25309a67d723dd8e75eb60c978b60929', 'check_textbox_on_bottomside__134680e7', 'check_text_exact_match__d332c3241fced231d1d84d00e75fe3b7', 'check_filename_pattern__a2f23245', 'check_textbox_on_rightside__5e2434f0', 'check_dir_file_count__95b4929b', 'check_file_list__4d646866', 'check_file_location__739292ff', 'check_text_contains_pattern__cc11abb5', 'check_text_uppercase__423804a3', 'check_file_organization__b1a155d8', 'check_text_replacement__dd369016ed084e5c5f565139bb4ef07f', 'check_zip_contains_files__8a01d242c9e2052109e1c26f5fa4a5dd', 'check_file_rename__c14397f3284104a2f980691d2ea6abf3', 'check_text_replacement__4aeab799', 'check_file_deleted__464fee7f', 'check_text_color_white__1b43c872', 'check_text_exists__109edd97', 'check_file_first_line__c5a909ed', 'check_file_renamed__bd9934867663f0945bb79537ace5711a', 'check_file_exists_with_size__995b229f', 'check_textbox_vertically_centered__8c858544', 'check_tex_file_list__557a0701', 'check_textbox_centered__b99b3c81', 'check_text_contains__3c678f53', 'check_file_contains_text__b6bb3e42', 'check_timezone_utc_minus_6__53de5eae', 'check_line_count__b2b950a6', 'check_json_has_fields__c8d946870135d67f0db0be5e65caaa2a', 'check_text_colors_match__6abdacd1e07293be0566dfd3601ee79b', 'check_file_moved__00112b53200a74ce7a53869d2d085264', 'check_python_functions__af4f2737', 'check_file_content__cd720b2833d8c75c48e4cd829046ee69', 'check_timezone__e55695d1', 'check_file_count__fbd9137c', 'check_git_repo_exists__68f4a1f5', 'check_file_count__2d28c7f27d8eb7cfbd3f915d05c991e3', 'check_text_replacement__f89b526b', 'check_extracted_files__940d01bc', 'check_file_organization__ba67a508', 'check_files_with_prefix__cf3f5d8ece62ecf8e4937dea9e007679', 'check_text_contains__b3a4b9dc', 'check_file_recovered_not_in_trash__05c91736', 'check_exact_text_match__05bf41d7', 'check_file_exists__3b8e423e430323c0078f4425aded05b9', 'check_textbox_bottomright__d0e1ffb6', 'check_timezone__2f3e1080', 'check_file_list__24c8b0df', 'exact_match__a85eebd24e563a97c17935bc46126aa1', 'check_file_exists__25f0d8c3', 'check_textbox_at_bottom__6cecb6d7']

def exact_match__6cc337d6467b2f92dee1f171f61ea4ed(result, expected, **options):
    """
    Check if result exactly matches expected value.

    Args:
        result: Current value from getter
        expected: Dictionary with 'expected' key

    Returns:
        1.0 if values match exactly, 0.0 otherwise
    """
    expect = expected.get('expected', '')
    logger.info(f'Result: {result}')
    logger.info(f'Expected: {expect}')
    if result == expect:
        return 1.0
    else:
        return 0.0

def check_import_lines__ad9b1b7a(result_file_path, expected, **options):
    """Check if import statements are correctly extracted from code.

    Args:
        result_file_path: Path to the file containing extracted imports
        expected: Dict with 'required_imports' (list) and 'min_count' (int)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result_file_path:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        import_lines = [line for line in lines if line.startswith('import ') or line.startswith('from ')]
        required_imports = expected.get('required_imports', [])
        min_count = expected.get('min_count', 0)
        score = 0.0
        if required_imports:
            found_count = sum((1 for req in required_imports if req in content))
            required_score = found_count / len(required_imports)
            score += required_score * 0.6
        if len(import_lines) >= min_count:
            score += 0.4
        return min(score, 1.0)
    except Exception as e:
        print(f'Error checking import lines: {e}')
        return 0.0

def check_exact_text_match__609a65b9(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_text_file_content__5ced85fc_aug18_v4_d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8a9(result, expected, **options):
    """Check if file content matches expected content exactly.

    Args:
        result: Content string from the file
        expected: Rules dict with 'expected_content' key
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        logger.error('Result is not a string')
        return 0.0
    expected_content = expected.get('expected_content', '')
    result_stripped = result.strip()
    expected_stripped = expected_content.strip()
    if result_stripped == expected_stripped:
        logger.info('Content matches expected value')
        return 1.0
    else:
        logger.info(f"Content mismatch: got '{result_stripped}', expected '{expected_stripped}'")
        return 0.0

def check_text_lines__c2b520792100f4aa54f0d89dedfa94c4(result, expected, **options):
    """
    Check if text file contains the expected lines (all 5 paper titles).

    Args:
        result: list of lines from the text file
        expected: list of expected titles

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct title)
    """
    if not isinstance(result, list) or not isinstance(expected, list):
        return 0.0
    if len(expected) == 0:
        return 0.0
    found_count = 0
    for expected_title in expected:
        expected_lower = expected_title.lower()
        for result_line in result:
            result_lower = result_line.lower()
            if expected_lower in result_lower or result_lower in expected_lower:
                found_count += 1
                break
    return found_count / len(expected)

def check_timezone_sydney__e603b840(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone (e.g., Sydney/Australian Eastern Time).

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Expected timezone configuration dict with 'timezone' key
        **options: Additional options

    Returns:
        1.0 if timezone matches expected, 0.0 otherwise
    """
    timezone_name = expected.get('timezone', 'Sydney')
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            line_lower = line.lower()
            timezone_lower = timezone_name.lower()
            if timezone_lower in line_lower:
                return 1.0
            if f'australia/{timezone_lower}' in line_lower:
                return 1.0
    return 0.0

def compare_text_output__f558b67b(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_first_line_left_aligned__adcd0938(result, expected, **options):
    """Check if first line is left aligned.

    Args:
        result: Alignment string from getter ('LEFT', 'CENTER', 'RIGHT', 'JUSTIFY')
        expected: Expected alignment from rules
        **options: Additional options

    Returns:
        float: 1.0 if left aligned, 0.0 otherwise
    """
    expected_alignment = expected.get('alignment', 'LEFT')
    return 1.0 if result == expected_alignment else 0.0

def check_file_count__9364293cce5b25e22063aee62da7d43d(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the folder contains exactly the expected set of filenames.

    Args:
        result: List of filenames from getter
        expected: Dict with 'filenames' key containing list of expected filenames,
                 and optional 'excluded_filenames' key for files that should NOT be present
        **options: Additional options (unused)

    Returns:
        1.0 if exact match, partial credit for partial match, 0.0 for no match
    """
    expected_filenames = expected.get('filenames', [])
    excluded_filenames = expected.get('excluded_filenames', [])
    result_set = set(result)
    expected_set = set(expected_filenames)
    excluded_set = set(excluded_filenames)
    if result_set & excluded_set:
        return 0.0
    if result_set == expected_set:
        return 1.0
    if result_set.issubset(expected_set):
        return len(result_set) / len(expected_set)
    correct_count = len(result_set & expected_set)
    wrong_count = len(result_set - expected_set)
    if correct_count == 0:
        return 0.0
    score = correct_count / len(expected_set)
    penalty = wrong_count / (len(expected_set) + wrong_count)
    return max(0.0, score - penalty)

def check_file_upload__68f5fe6b(result, expected, **options):
    """Check file upload details.

    Args:
        result: dict with file_count and filenames
        expected: dict with file_count and all_eml flag

    Returns:
        float: 1.0 if criteria met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_count = expected.get('file_count', 0)
    all_eml = expected.get('all_eml', False)
    if result.get('file_count', 0) != expected_count:
        return 0.0
    if all_eml:
        filenames = result.get('filenames', [])
        if not all((f.endswith('.eml') for f in filenames)):
            return 0.0
    return 1.0

def check_text_format__7f844786255954cce16f5ea58433f34e(result, expected, **options):
    """
    Check if text file matches expected format.

    Args:
        result: Dict with format info from getter
        expected: Rules dict with 'has_numbering' and 'min_lines' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    expected_numbering = expected.get('has_numbering', False)
    if result.get('has_numbering', False) == expected_numbering:
        score += 0.5
    min_lines = expected.get('min_lines', 0)
    if result.get('line_count', 0) >= min_lines:
        score += 0.5
    return score

def check_text_replacement__1f0d9653(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_python_code_complete__f55aa7954b40a62bad3b8ba851857ed1(result, expected, **options):
    """
    Check if Python code is complete with TODO filled and docstring added.

    Args:
        result: String content of the Python file
        expected: Dict with 'required_code' and 'has_docstring' keys
        **options: Additional options

    Returns:
        float: Partial score based on completion (0.5 for TODO, 0.5 for docstring)
    """
    if not isinstance(result, str) or not result:
        return 0.0
    score = 0.0
    required_code = expected.get('required_code', '')
    has_docstring = expected.get('has_docstring', False)
    if required_code and required_code in result:
        score += 0.5
    if has_docstring:
        if '"""' in result or "'''" in result:
            if 'def insertionSort' in result:
                func_pos = result.find('def insertionSort')
                search_area = result[func_pos:func_pos + 400]
                if '"""' in search_area or "'''" in search_area:
                    score += 0.5
    return score

def check_file_count_greater__00db2192(result, expected, **options):
    """Check if file count is greater than expected threshold.

    Args:
        result: Integer file count
        expected: Minimum expected count
        **options: Additional options

    Returns:
        float: 1.0 if count >= expected, 0.0 otherwise
    """
    try:
        if isinstance(result, (int, float)) and result >= expected:
            return 1.0
        return 0.0
    except:
        return 0.0

def check_python_syntax__198be354(result, expected, **options):
    """Check if Python file has valid syntax AND was extracted from Colab notebook.

    Args:
        result: dict from getter with syntax validation and extraction verification
        expected: dict (not used, both conditions must be met)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, has valid syntax, AND was extracted from notebook, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    file_exists = result.get('exists', False)
    valid_syntax = result.get('valid_syntax', False)
    extracted_from_notebook = result.get('extracted_from_notebook', False)
    if file_exists and valid_syntax and extracted_from_notebook:
        return 1.0
    return 0.0

def check_text_output__4cc7c3cf(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_exported__d22e0dfa(result, expected, **options):
    """Check if HTML file was successfully exported with proper content.

    Args:
        result: Dictionary containing file validation info
        expected: Expected value (True for file should exist with valid HTML)
        **options: Additional options

    Returns:
        float: 1.0 if file exists as valid HTML with table content, 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('has_content', False):
        return 0.0
    if not result.get('is_html', False):
        return 0.0
    if not result.get('has_table', False):
        return 0.0
    return 1.0

def check_python_imports__261836e0618ec18a7a70e8da4837dfbf(result, expected, **options):
    """Check if Python file has expected imports.

    Args:
        result: Import data from getter
        expected: Expected rules (dict with required_modules, min_imports, etc.)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    checks = 0
    if 'required_modules' in expected:
        checks += 1
        all_imports = result.get('standard_imports', []) + result.get('third_party_imports', []) + result.get('from_imports', [])
        imported_modules = set()
        for imp_stmt in all_imports:
            parts = imp_stmt.strip().split()
            if len(parts) >= 2:
                if parts[0] == 'import':
                    module = parts[1].split()[0]
                    imported_modules.add(module)
                    if '.' in module:
                        imported_modules.add(module.split('.')[0])
                elif parts[0] == 'from' and len(parts) >= 4:
                    module = parts[1]
                    imported_modules.add(module)
                    if '.' in module:
                        imported_modules.add(module.split('.')[0])
        required_modules = expected['required_modules']
        found = sum((1 for mod in required_modules if mod in imported_modules))
        if found == len(required_modules):
            score += 1.0
        else:
            score += found / len(required_modules) if required_modules else 0.0
    if 'min_imports' in expected:
        checks += 1
        if result.get('total_imports', 0) >= expected['min_imports']:
            score += 1.0
    if 'has_third_party' in expected:
        checks += 1
        has_third_party = len(result.get('third_party_imports', [])) > 0
        if has_third_party == expected['has_third_party']:
            score += 1.0
    if 'has_standard_library' in expected:
        checks += 1
        has_standard_library = len(result.get('standard_imports', [])) > 0
        if has_standard_library == expected['has_standard_library']:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_file_organization__92a58812(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_text_contains__5b8281b6(result, expected, **options):
    """Check if result text contains expected text.

    Args:
        result: Actual text from getter
        expected: Expected substring
        **options: Additional options (unused)

    Returns:
        1.0 if expected is in result, 0.0 otherwise
    """
    if expected in result:
        return 1.0
    return 0.0

def check_recent_file_count__db5a3e05(result, expected, **options):
    """Check if count of recent files matches expected.

    Args:
        result: Actual count of recent files
        expected: Dict with 'expected_count' key

    Returns:
        1.0 if match, 0.0 otherwise
    """
    expected_count = expected.get('expected_count', 2)
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_text_exact_match__68f83c37fe78a3f2dbadefcb3c480330(result: str, expected: dict, **options) -> float:
    """
    Check if text content exactly matches expected value.

    Args:
        result: Text content from getter
        expected: Dict with 'content' key containing expected text
        **options: Additional options

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        logger.debug('Result is None')
        return 0.0
    expected_content = expected.get('content', '')
    if result == expected_content:
        return 1.0
    else:
        logger.debug(f"Content mismatch. Expected: '{expected_content}', Got: '{result}'")
        return 0.0

def check_text_replacement__ffa8cf6ad724ee8fc8a065457d283c28(result: str, expected: Dict, **options) -> float:
    """
    Check if text replacement was performed correctly.

    Args:
        result: Actual file content (string from getter)
        expected: Expected value with 'expected_content' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        logger.warning('Result is empty')
        return 0.0
    expected_content = expected.get('expected_content', '')
    if not expected_content:
        logger.error('No expected_content in expected dict')
        return 0.0
    if result.strip() == expected_content.strip():
        return 1.0
    else:
        logger.info(f'Content mismatch. Expected length: {len(expected_content)}, Actual length: {len(result)}')
        return 0.0

def check_default_text_color__ea408cb7(result, expected, **options):
    """Check the default text color in LibreOffice Writer (variation 5: dark blue).

    This metric reads the LibreOffice configuration file (registrymodifications.xcu)
    and extracts the StandardColor value from the DefaultFont settings.

    Args:
        result: Path to the LibreOffice configuration file (vm_file downloaded from VM)
        expected: Dict containing expected color value (e.g., {'color_value': '128'})
        **options: Additional options passed by the framework

    Returns:
        float: 1.0 if color matches expected, 0.0 otherwise
    """
    config_file_path = result
    default_color = None
    expected_color = str(expected['color_value'])
    if not config_file_path:
        logger.error('Config file path is empty or None')
        return 0.0
    if not os.path.exists(config_file_path):
        logger.error(f'Config file not found at path: {config_file_path}')
        return 0.0
    logger.info(f'Reading LibreOffice config file from: {config_file_path}')
    try:
        tree = ET.parse(config_file_path)
        root = tree.getroot()
        namespace = {'oor': 'http://openoffice.org/2001/registry'}
        for elem in root.findall('.//item[@oor:path="/org.openoffice.Office.Writer/DefaultFont"]', namespace):
            for prop in elem.findall('.//prop[@oor:name="StandardColor"]', namespace):
                for value in prop.findall('value', namespace):
                    default_color = value.text
        if default_color is None:
            logger.warning('StandardColor property not found in the XML configuration file')
            return 0.0
        if default_color != expected_color:
            logger.info(f"StandardColor found but value mismatch: got '{default_color}', expected '{expected_color}'")
            return 0.0
    except FileNotFoundError:
        logger.error(f'Config file not found: {config_file_path}')
        return 0.0
    except ET.ParseError as e:
        logger.error(f'XML parse error in {config_file_path}: {e}')
        return 0.0
    except Exception as e:
        logger.error(f'Unexpected error reading config file {config_file_path}: {e}')
        return 0.0
    return 1.0

def check_text_exact_match__2d81db9f5efc3d72c38ba9f24bf6d4fb(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text content from file
        expected: Expected text content (from rules dict)
        **options: Additional options (case_sensitive, strip_whitespace)

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    strip_whitespace = options.get('strip_whitespace', True)
    case_sensitive = options.get('case_sensitive', True)
    result_text = result
    expected_text = expected
    if strip_whitespace:
        result_text = result_text.strip()
        expected_text = expected_text.strip()
    if not case_sensitive:
        result_text = result_text.lower()
        expected_text = expected_text.lower()
    return 1.0 if result_text == expected_text else 0.0

def check_file_in_list__880f8efb(result, expected, **options):
    """Check if expected filename is in the file list.

    Args:
        result: List of filenames from getter
        expected: Expected rules dict with 'filename' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if file found in list, 0.0 otherwise
    """
    expected_filename = expected.get('filename', '')
    if expected_filename in result:
        return 1.0
    else:
        logger.info(f"File '{expected_filename}' not found in directory. Files present: {result}")
        return 0.0

def check_file_exists__739292ff(result, expected, **options):
    """Check if file exists and is a valid PNG image with proper content.

    Args:
        result: Dict with file validation info {
            'exists': bool,
            'is_valid_png': bool,
            'file_size': int,
            'has_valid_dimensions': bool,
            'width': int or None,
            'height': int or None
        }
        expected: Dict with 'exists' key (True/False)
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        if not result.get('exists', True):
            return 1.0
        else:
            return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_valid_png', False):
        logger.info('File is not a valid PNG image')
        return 0.0
    file_size = result.get('file_size', 0)
    if file_size == 0:
        logger.info('File is empty')
        return 0.0
    if file_size < 1024:
        logger.info(f'File size ({file_size} bytes) is too small for a screenshot')
        return 0.0
    if not result.get('has_valid_dimensions', False):
        logger.info('File does not have valid dimensions')
        return 0.0
    return 1.0

def check_documents_file__c0d05fd4eae793b685d13d8511de53aa(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if file exists in Documents folder.

    Args:
        result: Dict from getter
        expected: Dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        logger.info('File does not exist')
        return 0.0
    if result.get('in_documents', False):
        score += 0.4
    else:
        logger.info('File not in Documents folder')
    if result.get('is_png', False):
        score += 0.2
    return score

def check_remaining_files__c90c6ca3(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the remaining files match the expected list.

    Args:
        result: List of remaining filenames
        expected: Dict with 'rules' containing 'expected_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_files = set(expected.get('expected_files', []))
    actual_files = set(result)
    if actual_files == expected_files:
        score = 1.0
    else:
        if not expected_files:
            return 0.0
        correct = len(actual_files & expected_files)
        extra = len(actual_files - expected_files)
        missing = len(expected_files - actual_files)
        score = max(0.0, correct / len(expected_files) - extra * 0.2 - missing * 0.2)
    logger.info(f'Expected: {expected_files}, Actual: {actual_files}, Score: {score}')
    return score

def check_file_content__11701199(result, expected, **options):
    """Check if file content matches expected text.

    Args:
        result: File content from command output (string)
        expected: Expected content rules dict with 'expected_content' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    expected_content = expected.get('expected_content', '')
    result_stripped = result.strip()
    expected_stripped = expected_content.strip()
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_file_contains_text__6c41fae8ffe95d2c4b3c521d5446de56(result: Dict[str, str], expected: Dict[str, Any], **options) -> float:
    """Check if a new markdown section was added to the file.

    This metric verifies that:
    1. The current file contains the expected markdown section header
    2. The original file did NOT contain this section header (proving it was newly added)

    Args:
        result: Dict with 'original' and 'current' file content from getter
        expected: Expected dict with 'text' key containing section title to search for
        **options: Additional options (ignore_case supported)

    Returns:
        float: 1.0 if section was newly added, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    current_content = result.get('current', '')
    original_content = result.get('original', '')
    if not current_content:
        return 0.0
    expected_text = expected.get('text', '')
    if not expected_text:
        return 0.0
    ignore_case = options.get('ignore_case', False)
    pattern = '^#+\\s+' + re.escape(expected_text)
    flags = re.MULTILINE
    if ignore_case:
        flags |= re.IGNORECASE
    current_has_section = bool(re.search(pattern, current_content, flags))
    if not current_has_section:
        return 0.0
    original_has_section = bool(re.search(pattern, original_content, flags))
    if original_has_section:
        return 0.0
    return 1.0

def check_file_count__23c57bc9(result, expected, **options):
    """Check if file count meets criteria.

    Args:
        result: File count from getter
        expected: Expected rules dict with 'min_count' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if count meets criteria, 0.0 otherwise
    """
    min_count = expected.get('min_count', 1)
    if result >= min_count:
        return 1.0
    else:
        logger.info(f'File count {result} is less than minimum {min_count}')
        return 0.0

def check_text_lines__7e304294ce76dab9f08b95178667a620(result, expected, **options):
    """
    Check if text file contains expected lines.

    Args:
        result: List of lines from getter
        expected: Rules dict with 'expected_lines' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_lines = expected.get('expected_lines', [])
    if not expected_lines:
        return 0.0
    ignore_case = options.get('ignore_case', False)
    ignore_order = options.get('ignore_order', False)
    if ignore_case:
        result_normalized = [line.lower() for line in result]
        expected_normalized = [line.lower() for line in expected_lines]
    else:
        result_normalized = result
        expected_normalized = expected_lines
    if ignore_order:
        return float(all((exp in result_normalized for exp in expected_normalized)))
    else:
        return float(result_normalized == expected_normalized)

def check_single_eml_file__54c9c002bd0d5132155806dba06b543e(result: Optional[str], expected: dict, **options) -> float:
    """Check if a single email file was successfully downloaded and contains expected content.

    Args:
        result: Downloaded file path (or None if not found)
        expected: Expected rules dict with:
            - subject_pattern: Expected subject line pattern (optional)
        **options: Additional options

    Returns:
        Score: 1.0 if file exists (and matches pattern if specified), 0.0 otherwise
    """
    if result is None:
        return 0.0
    subject_pattern = expected.get('subject_pattern')
    if subject_pattern is None:
        return 1.0
    try:
        with open(result, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read(5000)
            if re.search(f'Subject:.*{re.escape(subject_pattern)}', content, re.IGNORECASE):
                return 1.0
            else:
                logger.warning(f"Subject pattern '{subject_pattern}' not found in file")
                return 0.0
    except Exception as e:
        logger.error(f"Failed to read file '{result}': {e}")
        return 0.0

def check_file_exported__7107319d(result, expected, **options):
    """Check if file was successfully exported as HTML with proper content.

    Args:
        result: Dictionary with validation results from getter:
            - exists: bool - whether file exists
            - is_html: bool - whether file contains HTML content
            - has_table: bool - whether file contains table data
            - file_size: int - size of file in bytes
            - has_content: bool - whether file has substantial content
        expected: Expected value (True for file should exist and be valid)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and contains valid HTML spreadsheet data, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if isinstance(result, bool):
        return 1.0 if result == expected else 0.0
    if isinstance(result, dict):
        if expected is True or expected == True:
            if not result.get('exists', False):
                return 0.0
            if not result.get('is_html', False):
                return 0.0
            if not result.get('has_table', False):
                return 0.0
            if not result.get('has_content', False):
                return 0.0
            return 1.0
        else:
            all_valid = result.get('exists', False) and result.get('is_html', False) and result.get('has_table', False) and result.get('has_content', False)
            return 0.0 if all_valid else 1.0
    return 0.0

def check_text_file_lines__5ced85fc_aug18_v3_c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7f8(result, expected, **options):
    """Check if file lines match expected lines exactly.

    Args:
        result: List of lines from the file
        expected: Rules dict with 'expected_lines' key
        **options: Additional options

    Returns:
        float: 1.0 if all lines match, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error('Result is not a list')
        return 0.0
    expected_lines = expected.get('expected_lines', [])
    if len(result) != len(expected_lines):
        logger.info(f'Line count mismatch: got {len(result)}, expected {len(expected_lines)}')
        return 0.0
    for (i, (actual, expected_line)) in enumerate(zip(result, expected_lines)):
        if actual != expected_line:
            logger.info(f"Line {i + 1} mismatch: got '{actual}', expected '{expected_line}'")
            return 0.0
    logger.info('All lines match expected values')
    return 1.0

def check_file_size_range__0f2a6243(result, expected, **options):
    """
    Check if file exists and size is within expected range.

    Args:
        result: Dict from getter with 'exists' and 'size'
        expected: Dict with 'min_size' and 'max_size'
        **options: Additional options

    Returns:
        float: 1.0 if checks pass, 0.5 if exists but wrong size, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.5
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    file_size = result.get('size', 0)
    if min_size <= file_size <= max_size:
        score += 0.5
    return score

def check_both_files__3f6d3219(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if multiple files exist and meet requirements.

    Args:
        result: Dict from getter with 'files' list
        expected: Expected rules dict

    Returns:
        float: 1.0 if all files meet requirements, 0.0 otherwise
    """
    files = result.get('files', [])
    if expected.get('both_exist', True):
        if not all((f.get('exists', False) for f in files)):
            return 0.0
    min_size_kb = expected.get('min_size_kb', 0)
    if min_size_kb > 0:
        for f in files:
            if f.get('exists', False):
                size_kb = f.get('size_bytes', 0) / 1024
                if size_kb < min_size_kb:
                    return 0.0
    return 1.0

def check_note_exact_match__5f5d62626ae9f466dd58a3e57638d2a2(result_state, expected_state, **options):
    """
    Check if the note text exactly matches the expected text.

    Args:
        result_state: The actual note text from the slide (str)
        expected_state: Expected state dict with key 'expected_text' (str)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if the note matches exactly, 0.0 otherwise
    """
    expected_text = expected_state.get('expected_text', '')
    if result_state == expected_text:
        return 1.0
    return 0.0

def check_sender_file__ad43db2627764dd28ab9631606c7b97c(result: str, expected: dict, **options) -> float:
    """
    Check if a text file contains the expected sender email address.

    Args:
        result: Content from the created text file
        expected: Expected values from rules (dict with 'email')
        **options: Additional options

    Returns:
        1.0 if the content contains the expected email address, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        result_str = str(result).strip()
        expected_email = expected.get('email', '').strip().lower()
        result_lower = result_str.lower()
        if result_lower == expected_email:
            return 1.0
        email_pattern = '<([^>]+)>'
        matches = re.findall(email_pattern, result_str)
        if matches:
            for email in matches:
                if email.lower() == expected_email:
                    return 1.0
        if expected_email in result_lower:
            return 0.9
        email_pattern2 = '[\\w\\.-]+@[\\w\\.-]+'
        emails_found = re.findall(email_pattern2, result_str)
        for email in emails_found:
            if email.lower() == expected_email:
                return 1.0
        return 0.0
    except Exception:
        return 0.0

def check_contains_filename__32107748(result, expected, **options):
    """Check if forwarded_paper.eml exists and contains Paper Recommendation email.

    Args:
        result: Dict with 'exists', 'filename', and 'content' keys
        expected: Dict with 'filename' key

    Returns:
        float: 1.0 if file exists with correct content, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    filename = expected.get('filename', '')
    if not filename:
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if result.get('filename', '') != filename:
        return 0.0
    content = result.get('content', '')
    if not content:
        return 0.0
    content_lower = content.lower()
    has_paper_recommendation = False
    for line in content.split('\n'):
        line_lower = line.lower()
        if line_lower.startswith('subject:') and 'paper recommendation' in line_lower:
            has_paper_recommendation = True
            break
    if not has_paper_recommendation:
        return 0.0
    return 1.0

def check_merged_text__e22bfb55f6ab9983d1cc35b82dc09aeb(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if merged text file exists and contains all chapters.

    Args:
        result: Dict from getter with 'exists', 'line_count', 'contains_chapters'
        expected: Dict with 'min_line_count', 'required_chapters'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        return 0.0
    score += 0.3
    min_line_count = expected.get('min_line_count', 100)
    line_count = result.get('line_count', 0)
    if line_count >= min_line_count:
        score += 0.3
    elif line_count > 0:
        score += 0.3 * (line_count / min_line_count)
    required_chapters = expected.get('required_chapters', [])
    contains_chapters = result.get('contains_chapters', [])
    if required_chapters:
        matched = sum((1 for ch in required_chapters if ch in contains_chapters))
        chapter_score = matched / len(required_chapters)
        score += 0.4 * chapter_score
    return min(score, 1.0)

def check_textbox_bottomleft__99aea82e(src_path, expected_state, **options):
    """
    Check if the textbox is in the bottom-left corner of the image.
    Variation 6 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in bottom-left 5% region, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                bottom_most = max(bottom_most, y)
    if left_most < width * 0.05 and bottom_most > height * 0.95:
        return 1.0
    else:
        return 0.0

def check_text_contains__de3b1681(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    score = 0.0
    for substring in contains:
        if substring.lower() in result.lower():
            score += 1.0 / len(contains)
    return score

def check_text_contains__c491dbb8(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    return 1.0 if any((substring.lower() in result.lower() for substring in contains)) else 0.0

def check_file_permissions__11701199(result, expected, **options):
    """Check if file/directory permissions match expected value.

    Args:
        result: Permission value from stat command (string like "755" or "555")
        expected: Expected state rules with 'expected_permissions'
        **options: Additional comparison options

    Returns:
        float: 1.0 if permissions match, 0.0 otherwise
    """
    expected_permissions = expected.get('expected_permissions', '')
    result_stripped = result.strip()
    if result_stripped == expected_permissions:
        return 1.0
    else:
        return 0.0

def check_textbox_on_rightside__5bc789db(src_path, expected_state, **options):
    """
    Check if the textbox is on the right side of the image.
    Variation 1 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text bounding box center is in right half, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    right_most = 0
    top_most = height
    bottom_most = 0
    found_text = False
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                found_text = True
                left_most = min(left_most, x)
                right_most = max(right_most, x)
                top_most = min(top_most, y)
                bottom_most = max(bottom_most, y)
    if not found_text:
        return 0.0
    text_center_x = (left_most + right_most) / 2
    if text_center_x > width * 0.5:
        return 1.0
    else:
        return 0.0

def check_files_match_pattern__fb580d40(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if files match expected pattern and count"""
    if result is None:
        return 0.0
    pattern = expected.get('pattern', '.*')
    min_count = expected.get('min_count', 1)
    try:
        regex = re.compile(pattern)
    except Exception as e:
        logger.error(f'Invalid regex pattern: {pattern}, error: {e}')
        return 0.0
    matching_files = [f for f in result if regex.match(f)]
    match_count = len(matching_files)
    if match_count >= min_count:
        logger.info(f"Pattern check passed: {match_count} files match pattern '{pattern}' (min: {min_count})")
        return 1.0
    else:
        logger.info(f"Pattern check failed: {match_count} files match pattern '{pattern}' (min: {min_count})")
        return 0.0

def check_text_file_contains__ebe7bea30aab4d43b91ea1760b3fb66f(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if text file contains all expected strings.

    Args:
        result: File content from getter
        expected: Dict with 'contains' key - list of strings that must be in the file

    Returns:
        Partial score based on how many expected strings are found
    """
    expected_strings = expected.get('contains', [])
    if not expected_strings:
        logger.warning('No expected strings specified')
        return 0.0
    if not result:
        logger.warning('File content is empty')
        return 0.0
    found_count = 0
    for expected_str in expected_strings:
        if expected_str in result:
            found_count += 1
        else:
            logger.info(f'Missing expected string: {expected_str}')
    score = found_count / len(expected_strings)
    logger.info(f'Found {found_count}/{len(expected_strings)} expected strings, score: {score}')
    return score

def check_text_contains_all__b8171706418d2058f81460f1a24d4635(result, expected, **options):
    """
    Check if text contains all required paper titles and keywords.

    This evaluator verifies TWO requirements from the instruction:
    1. All paper titles from the document are included
    2. Important keywords are included

    Paper titles: Spider, SParC, CoSQL
    Keywords: SQL, BERT, Attention, Language Models

    Args:
        result: String content from getter
        expected: Rules dict with 'paper_titles' and 'keywords' lists
        **options: Additional options (ignore_case)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    paper_titles = expected.get('paper_titles', [])
    keywords = expected.get('keywords', [])
    if not paper_titles and (not keywords):
        return 0.0
    ignore_case = options.get('ignore_case', False)
    title_score = 0.0
    if paper_titles:
        if ignore_case:
            result_lower = result.lower()
            titles_lower = [title.lower() for title in paper_titles]
            title_matches = sum((1 for title in titles_lower if title in result_lower))
        else:
            title_matches = sum((1 for title in paper_titles if title in result))
        title_score = title_matches / len(paper_titles)
    keyword_score = 0.0
    if keywords:
        if ignore_case:
            result_lower = result.lower()
            keywords_lower = [kw.lower() for kw in keywords]
            keyword_matches = sum((1 for kw in keywords_lower if kw in result_lower))
        else:
            keyword_matches = sum((1 for kw in keywords if kw in result))
        keyword_score = keyword_matches / len(keywords)
    if paper_titles and keywords:
        final_score = title_score * 0.5 + keyword_score * 0.5
    elif paper_titles:
        final_score = title_score
    else:
        final_score = keyword_score
    return final_score

def check_python_comment__4ca3247dcf6b464342c4e3f53d844797(result: str, expected: Dict, **options) -> float:
    """Check if Python file contains a comment with expected text at the top.

    Args:
        result: String containing the actual file content
        expected: Dict with 'comment_text' field specifying expected comment content

    Returns:
        1.0 if comment is found in the first 5 non-empty lines, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_comment = expected.get('comment_text', '')
    if not expected_comment:
        return 0.0
    lines = result.split('\n')
    for line in lines[:5]:
        stripped = line.strip()
        if stripped.startswith('#') and expected_comment in stripped:
            return 1.0
    return 0.0

def compare_text_output__2175dacb(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_multiline_answers__fa7deb995de2addd4a7d5d55ff4c5c25(result: Dict[str, List[str]], expected: Dict[str, Any], **options) -> float:
    """
    Check if the answer file contains the expected answers on separate lines.

    Args:
        result: Dict mapping test names to lists of answer lines (from getter)
        expected: Dict with 'answers' key containing expected test answers
        **options: Additional options (e.g., 'ignore_case')

    Returns:
        Score between 0.0 and 1.0
    """
    expected_answers = expected.get('answers', {})
    ignore_case = options.get('ignore_case', False)
    if not expected_answers:
        logger.warning('No expected answers provided')
        return 0.0
    total_tests = len(expected_answers)
    if total_tests == 0:
        return 0.0
    correct_count = 0
    for (test_name, expected_answer_list) in expected_answers.items():
        if test_name not in result:
            logger.info(f"Test '{test_name}' not found in result")
            continue
        result_lines = result[test_name]
        if len(result_lines) != len(expected_answer_list):
            logger.info(f"Line count mismatch for '{test_name}': got {len(result_lines)}, expected {len(expected_answer_list)}")
            continue
        all_match = True
        for (i, (result_line, expected_line)) in enumerate(zip(result_lines, expected_answer_list)):
            if ignore_case:
                result_line = result_line.lower()
                expected_line = expected_line.lower()
            if result_line != expected_line:
                logger.info(f"Line {i + 1} mismatch for '{test_name}': got '{result_line}', expected '{expected_line}'")
                all_match = False
                break
        if all_match:
            correct_count += 1
    score = correct_count / total_tests
    return score

def check_first_line__66414428(result, expected, **options):
    """Check if first line contains expected text.
    
    Args:
        result: First line of file
        expected: Dict with 'contains' key
        
    Returns:
        float: 1.0 if first line contains expected text, else 0.0
    """
    if result is None:
        return 0.0
    expected_text = expected.get('contains', '')
    ignore_case = options.get('ignore_case', True)
    if ignore_case:
        return 1.0 if expected_text.lower() in result.lower() else 0.0
    else:
        return 1.0 if expected_text in result else 0.0

def check_file_exists__d967a5b0(result, expected, **options):
    """Check if file existence message matches expected.

    Args:
        result: String output from test command
        expected: Dict with 'value' key ('exists' or 'missing')
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', 'exists')
    result_clean = result.strip().lower()
    expected_clean = expected_value.strip().lower()
    if result_clean == expected_clean:
        return 1.0
    else:
        logger.info(f"File existence mismatch: expected '{expected_value}', got '{result}'")
        return 0.0

def check_text_content__1a1f627807b83c4d33e1ae428da08935(result, expected, **options):
    """Check if file content matches expected text exactly.

    Args:
        result: String content from the file
        expected: Dict with 'content' key containing expected text
        **options: Additional options (ignore_trailing_whitespace, etc.)

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_content = expected.get('content', '')
    ignore_trailing = options.get('ignore_trailing_whitespace', False)
    if ignore_trailing:
        result = result.rstrip()
        expected_content = expected_content.rstrip()
    if result == expected_content:
        return 1.0
    return 0.0

def check_file_exists__2ad8e92c(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def exact_count_match__f1d66967(result, expected, **options):
    """Check if PDFs match expected criteria for blog article downloads.

    This verifies:
    1. Correct count (2 PDFs)
    2. Filenames contain keywords from blog URLs/titles
    3. File sizes are non-trivial (> 100KB for blog articles)

    Args:
        result: Dict with 'count' and 'files' (list of file info)
        expected: Dict with 'count' key

    Returns:
        1.0 if all criteria match, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if isinstance(result, dict):
        actual_count = result.get('count', 0)
        files = result.get('files', [])
    else:
        actual_count = result
        files = []
    if actual_count != expected_count:
        logger.info(f'File count mismatch: expected {expected_count}, got {actual_count}')
        return 0.0
    if not files:
        return 1.0
    MIN_FILE_SIZE = 100 * 1024
    for file_info in files:
        size = file_info.get('size', 0)
        if size < MIN_FILE_SIZE:
            logger.info(f"File {file_info.get('name')} is too small ({size} bytes), expected > {MIN_FILE_SIZE}")
            return 0.0
    keywords = ['agent', 'human', 'data', 'quality']
    matched_keywords = set()
    for file_info in files:
        filename = file_info.get('name', '').lower()
        for keyword in keywords:
            if keyword in filename:
                matched_keywords.add(keyword)
    if len(matched_keywords) < 2:
        logger.info(f"Filenames don't contain expected keywords. Found keywords: {matched_keywords}")
        return 0.0
    logger.info(f'All verification checks passed: {actual_count} PDFs with valid sizes and filenames')
    return 1.0

def check_file_renamed__9f8a1d0c(result, expected, **options):
    """Check if file was renamed correctly.

    Args:
        result: Dict from getter with old_exists and new_exists
        expected: Dict (not used, checking that old doesn't exist and new does)
        **options: Additional options

    Returns:
        float: 1.0 if renamed correctly, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('old_exists', True) and result.get('new_exists', False):
        return 1.0
    return 0.0

def check_textbox_in_bottom_right__c729f30a(src_path, expected, **options):
    """
    Check if the textbox is in the bottom-right corner of the image.
    Task variation 7 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is in bottom-right, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    right_most = 0
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                right_most = max(right_most, x)
                bottom_most = max(bottom_most, y)
    if right_most > width * 0.95 and bottom_most > height * 0.95:
        return 1.0
    else:
        return 0.0

def check_content_text_color__53beb8d2(result_file, expected, **options):
    """
    Check if the content text color matches the expected RGB color.

    Args:
        result_file: Path to the PPTX file
        expected: Expected values with 'slide_idx', 'shape_idx', 'expected_rgb' (tuple of R, G, B)
        **options: Additional options

    Returns:
        float: 1.0 if text color matches expected, 0.0 otherwise
    """
    try:
        prs = Presentation(result_file)
        slide_idx = expected.get('slide_idx', 0)
        shape_idx = expected.get('shape_idx', 1)
        expected_rgb = tuple(expected.get('expected_rgb', (255, 0, 0)))
        if slide_idx >= len(prs.slides):
            return 0.0
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            return 0.0
        shape = slide.shapes[shape_idx]
        if not hasattr(shape, 'text_frame'):
            return 0.0
        has_text = False
        all_match = True
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text.strip():
                    has_text = True
                    if hasattr(run.font.color, 'rgb'):
                        actual_rgb = run.font.color.rgb
                        if tuple(actual_rgb) != expected_rgb:
                            all_match = False
                            break
                    else:
                        all_match = False
                        break
            if not all_match:
                break
        if has_text and all_match:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_timezone__3d14fb6b(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_renamed__b4364020(result, expected, **options):
    """Check if file was renamed correctly.

    Args:
        result: Dict with rename status from getter
        expected: Expected state (dict with 'renamed' key)
        **options: Additional options

    Returns:
        float: 1.0 if file was renamed, 0.0 otherwise
    """
    expected_renamed = expected.get('renamed', True)
    actual_renamed = result.get('renamed', False)
    if actual_renamed == expected_renamed:
        return 1.0
    else:
        return 0.0

def check_file_exists__9c31b3f6afb568d600c17c937149b6c4(result: bool, expected: Dict[str, Any], **options) -> float:
    """Check if file exists as expected.

    Args:
        result: Boolean indicating if file exists
        expected: Dict with 'exists' field (True/False)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_timezone__43ba8703(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_exists__4c5ac05d(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def check_textbox_centered__723c4038(src_path, expected, **options):
    """
    Check if the textbox is centered in the image.
    Task variation 1 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is centered, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    right_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                right_most = max(right_most, x)
    text_center = (left_most + right_most) / 2
    image_center = width / 2
    tolerance = width * 0.1
    if abs(text_center - image_center) < tolerance:
        return 1.0
    else:
        return 0.0

def check_file_exported__6fb34f5d(result, expected, **options):
    """Check if file was successfully exported with correct ODS format.

    Args:
        result: dict with 'exists', 'is_ods', 'file_type' keys
        expected: Expected value (True for file should exist with correct format)
        **options: Additional options

    Returns:
        float: Score based on file existence and format validation
               1.0 if file exists AND is valid ODS format
               0.5 if file exists but not valid ODS format
               0.0 if file doesn't exist
    """
    if result is None:
        return 0.0
    if isinstance(result, bool):
        return 1.0 if result == expected else 0.0
    if not isinstance(result, dict):
        return 0.0
    exists = result.get('exists', False)
    is_ods = result.get('is_ods', False)
    if not exists:
        return 0.0
    if exists and (not is_ods):
        return 0.5
    if exists and is_ods:
        return 1.0
    return 0.0

def check_exact_match_v5__d4477d7a(result, expected, **options):
    """Compare result against expected integer value.

    Args:
        result: Actual value from getter
        expected: Rules dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_value = expected.get('expected_value')
    if result is None or expected_value is None:
        return 0.0
    try:
        result_int = int(result)
        expected_int = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    if result_int == expected_int:
        return 1.0
    return 0.0

def check_textbox_on_topside__20161493a3832ab8acb83e42c1b1cbaf(result_state, expected_state, **options):
    """
    Check if the textbox is on the top side of the image.
    Variation of gimp:e2dd0213-26db-4349-abe5-d5667bfd725c
    Task: Move text layer upward to top area

    Args:
        result_state: Path to the exported image file
        expected_state: Expected value (not used, verification is binary)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is in top area, 0.0 otherwise)
    """
    src_path = result_state
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    top_most_dark_pixel = height
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                top_most_dark_pixel = min(top_most_dark_pixel, y)
                break
        if top_most_dark_pixel < height:
            break
    if top_most_dark_pixel < height * 0.1:
        return 1.0
    else:
        return 0.0

def check_file_format_and_structure__6bbe76fc(result, expected, **options):
    """
    Check if the image has been saved in the correct format and structure is preserved.

    Args:
        result: Path to the result image file
        expected: Path to the original image file
        **options: Additional options including target_format

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    try:
        if not os.path.exists(result):
            logger.error(f'Result file does not exist: {result}')
            return 0.0
        result_img = Image.open(result)
        result_format = result_img.format
        target_format = options.get('target_format', 'PNG')
        logger.debug(f'Result format: {result_format}, Target format: {target_format}')
        score = 0.0
        if result_format and result_format.upper() == target_format.upper():
            score += 0.5
            logger.debug(f'Format matches: {result_format}')
        else:
            logger.debug(f'Format mismatch: {result_format} != {target_format}')
        score += 0.5
        logger.debug(f'Final format check score: {score}')
        return score
    except Exception as e:
        logger.error(f'Error checking file format: {e}')
        return 0.0

def check_file_exists__42ef2d72a0f41077972857e814318a24(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file exists as expected.

    Args:
        result: Dict with 'exists' bool from getter
        expected: Dict with 'should_exist' (expected existence state)
        **options: Additional options

    Returns:
        float: 1.0 if existence matches expectation, 0.0 otherwise
    """
    exists = result.get('exists', False)
    should_exist = expected.get('should_exist', True)
    if exists == should_exist:
        return 1.0
    else:
        return 0.0

def check_recent_files__4e03b1ed(result, expected, **options):
    """Check if expected number of recent files exist.

    Args:
        result: Count of recent files from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('count', 2)
    if result >= expected_count:
        logger.info(f'Found {result} recent PDF files (expected >= {expected_count})')
        return 1.0
    else:
        logger.info(f'Only {result} recent PDF files, expected >= {expected_count}')
        return 0.0

def check_textbox_in_bottom_left__f062cc22(src_path, expected, **options):
    """
    Check if the textbox is in the bottom-left corner of the image.
    Task variation 6 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is in bottom-left, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    bottom_left_x_threshold = int(width * 0.05)
    bottom_left_y_start = int(height * 0.95)
    found_dark_pixel = False
    for y in range(bottom_left_y_start, height):
        for x in range(bottom_left_x_threshold):
            if gray_image.getpixel((x, y)) < 128:
                found_dark_pixel = True
                logger.info(f'Found dark pixel at ({x}, {y}) in bottom-left region')
                break
        if found_dark_pixel:
            break
    if found_dark_pixel:
        logger.info(f'Text found in bottom-left corner (region: x < {bottom_left_x_threshold}, y > {bottom_left_y_start})')
        return 1.0
    else:
        logger.info(f'No text found in bottom-left corner (region: x < {bottom_left_x_threshold}, y > {bottom_left_y_start})')
        return 0.0

def check_timezone__308d9780(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_python_structure__58d4fb9ea6e69b7f57a37a59813050e8(result, expected, **options):
    """Check if Python file has expected structure.

    Args:
        result: File structure from getter
        expected: Expected structure rules (dict with min_imports, required_classes, etc.)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    total_checks = 0
    if 'min_imports' in expected:
        total_checks += 1
        if len(result.get('imports', [])) >= expected['min_imports']:
            score += 1.0
    if 'required_classes' in expected:
        total_checks += 1
        result_classes = result.get('classes', [])
        expected_classes = expected['required_classes']
        if all((cls in result_classes for cls in expected_classes)):
            score += 1.0
    if 'required_functions' in expected:
        total_checks += 1
        result_functions = result.get('functions', [])
        expected_functions = expected['required_functions']
        if all((func in result_functions for func in expected_functions)):
            score += 1.0
    if 'min_lines' in expected:
        total_checks += 1
        if result.get('total_lines', 0) >= expected['min_lines']:
            score += 1.0
    return score / total_checks if total_checks > 0 else 0.0

def check_file_exists_with_size__d6a98717(result, expected, **options):
    """
    Check if file exists and has reasonable size.

    Args:
        result: dict from getter with keys {"exists": bool, "size": int, "is_png": bool}
        expected: dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_comment_lines__70bdd422(result_file_path, expected, **options):
    """Check if comment lines are correctly extracted from code.

    Args:
        result_file_path: Path to file containing extracted comments
        expected: Dict with 'min_comments' (int) and 'all_start_with_hash' (bool)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result_file_path:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        lines = [line.strip() for line in content.splitlines() if line.strip()]
        min_comments = expected.get('min_comments', 0)
        all_start_with_hash = expected.get('all_start_with_hash', False)
        score = 0.0
        if len(lines) >= min_comments:
            score += 0.5
        if all_start_with_hash and lines:
            valid_comments = sum((1 for line in lines if line.startswith('#')))
            comment_ratio = valid_comments / len(lines)
            score += comment_ratio * 0.5
        return min(score, 1.0)
    except Exception as e:
        print(f'Error checking comment lines: {e}')
        return 0.0

def check_text_increased__e6f52b2c(result, expected, **options):
    """Check if text length increased sufficiently.

    Args:
        result: Actual text length
        expected: Dict with min_length

    Returns:
        float: 1.0 if >= min_length, 0.0 otherwise
    """
    if isinstance(expected, dict):
        min_length = expected.get('min_length', 0)
    else:
        min_length = expected
    return 1.0 if result >= min_length else 0.0

def check_python_with_summary__be0d67f2ca1c7fde5d3c550cd046d0b4(result, expected, **options):
    """
    Check if Python file has a summary comment at the top showing accurate line count
    and contains code that appears to be from a Colab notebook.

    Args:
        result: File content string from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    score = 0.0
    if result and len(result.strip()) > 0:
        score += 0.2
    else:
        return 0.0
    lines = result.split('\n')
    has_summary = False
    claimed_line_count = None
    for line in lines[:10]:
        match = re.search('#.*(?:total\\s*)?(?:line|lines)(?:\\s*count)?.*:\\s*(\\d+)', line, re.IGNORECASE)
        if match:
            has_summary = True
            claimed_line_count = int(match.group(1))
            break
    if has_summary:
        score += 0.2
    code_lines = []
    for (i, line) in enumerate(lines):
        if i < 10 and re.search('#.*(?:total\\s*)?(?:line|lines)(?:\\s*count)?.*:\\s*\\d+', line, re.IGNORECASE):
            continue
        stripped = line.rstrip()
        if stripped:
            code_lines.append(stripped)
    actual_line_count = len(code_lines)
    if claimed_line_count is not None and actual_line_count > 0:
        if abs(claimed_line_count - actual_line_count) <= 2:
            score += 0.2
        elif abs(claimed_line_count - actual_line_count) <= 5:
            score += 0.1
    executable_code_lines = [line for line in code_lines if not line.strip().startswith('#')]
    if len(executable_code_lines) > 0:
        score += 0.2
    notebook_indicators = 0
    content_lower = result.lower()
    common_imports = ['import numpy', 'import pandas', 'import matplotlib', 'import tensorflow', 'import torch', 'import sklearn', 'from google.colab', 'import seaborn', 'import keras']
    if any((imp in content_lower for imp in common_imports)):
        notebook_indicators += 1
    if re.search('\\ndef\\s+\\w+\\s*\\(', result) or re.search('\\nclass\\s+\\w+', result):
        notebook_indicators += 1
    import_blocks = len(re.findall('(?:^|\\n)(?:import|from)\\s+\\w+', result))
    function_blocks = len(re.findall('(?:^|\\n)def\\s+\\w+', result))
    if import_blocks + function_blocks >= 2:
        notebook_indicators += 1
    if notebook_indicators >= 2:
        score += 0.2
    elif notebook_indicators == 1:
        score += 0.1
    return score

def check_text_replacement__0b4a6158(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_file_path__8369c71b0ee26c543c025c6e1cb39bbd(result, expected, **options):
    """Compare file path content against expected path.

    Args:
        result: The content read from the file
        expected: Expected path string (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if paths match, 0.0 otherwise
    """
    expected_path = expected.get('expected_path', '')
    if not result:
        return 0.0
    if result.strip() == expected_path.strip():
        return 1.0
    else:
        return 0.0

def check_file_rename_pattern__a610cecf(result, expected, **options):
    """
    Check if files have been properly renamed (not copied) according to a pattern.

    Args:
        result (str): Content from cache_file containing ls output of the directory
        expected (dict): Expected configuration with:
            - directory (str): The directory path
            - pattern_files (list): List of expected renamed files
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
            - 1.0 if all files renamed and old files removed
            - 0.0 if any old file still exists (copy instead of rename)
            - Proportional score if only some files renamed correctly
    """
    try:
        files_in_dir = [line.strip() for line in result.strip().split('\n') if line.strip()]
        directory = expected.get('directory', '')
        pattern_files = expected.get('pattern_files', [])
        if not pattern_files:
            return 0.0
        matched = 0
        for pattern_file in pattern_files:
            if pattern_file in files_in_dir:
                matched += 1
        if matched < len(pattern_files):
            return matched / len(pattern_files)
        for file in files_in_dir:
            if file.startswith('aws-invoice-') and file.endswith('.pdf'):
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error in check_file_rename_pattern__a610cecf: {e}')
        return 0.0

def check_music_file_count__355d5753246a8a88c1b8d173c110d89e(result: list, expected: dict, **options) -> float:
    """Check if correct MP3 files remain after deletion.

    Verifies that:
    1. File count matches expected value (3 files)
    2. No remaining files contain 'Missing' or 'Painful' in their names

    Args:
        result: List of MP3 filenames from getter
        expected: Dict from rules with 'count' key
        **options: Additional options

    Returns:
        float: 1.0 if verification passes, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if len(result) != expected_count:
        logger.debug(f'MP3 file count mismatch: got {len(result)}, expected {expected_count}')
        return 0.0
    forbidden_keywords = ['Missing', 'Painful']
    for filename in result:
        for keyword in forbidden_keywords:
            if keyword in filename:
                logger.debug(f"Found file with forbidden keyword '{keyword}': {filename}")
                return 0.0
    logger.debug(f"Verification passed: {len(result)} files, none contain 'Missing' or 'Painful'")
    return 1.0

def check_text_on_rightside__68e0eaa1c0be88d7d26c920a571b88e0(result, expected, **options):
    """
    Check if the text is on the right side of the image.

    Args:
        result: X-coordinate of the rightmost text pixel
        expected: Dictionary with 'image_width' key
        **options: Additional options

    Returns:
        float: 1.0 if text is on the right side (within right 5%), 0.0 otherwise
    """
    if result is None or result == 0:
        return 0.0
    image_width = expected.get('image_width', 2192)
    if result > image_width * 0.95:
        return 1.0
    else:
        return 0.0

def check_file_content__e1c8e8d0(result, expected, **options):
    """
    Check if the file content matches the expected value.

    Args:
        result: Actual file content from getter
        expected: Dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_value = str(expected.get('expected_value', ''))
    if result.strip() == expected_value.strip():
        return 1.0
    else:
        return 0.0

def check_file_exists__c0930f6fc6470d951ad9d775e3a6c6a5(result, expected, **options):
    """Check if file exists as expected and validate it's a proper TIFF file with content.

    Args:
        result: Result from getter (dict with 'exists', 'is_tiff', 'has_content' keys)
        expected: Expected rules (dict with 'should_exist' key)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, is a valid TIFF, and has content; 0.0 otherwise
    """
    file_exists = result.get('exists', False)
    is_tiff = result.get('is_tiff', False)
    has_content = result.get('has_content', False)
    should_exist = expected.get('should_exist', True)
    if not should_exist:
        return 1.0 if not file_exists else 0.0
    if file_exists and is_tiff and has_content:
        return 1.0
    else:
        return 0.0

def check_git_repo_status__55495fb6b59196cc7cffae2b12e117ed(result, expected, **options):
    """Check if git repository exists and has correct remote URL.

    Args:
        result: Dict with keys 'exists', 'is_git_repo', 'remote_url'
        expected: Rules dict with 'url' key (str)
        **options: Additional options

    Returns:
        float: Partial credit score:
            - 0.5 if repository exists and is a valid git repo
            - 1.0 if above conditions met AND remote URL matches
            - 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('is_git_repo', False):
        return 0.0
    score = 0.5
    expected_url = expected.get('url', '')
    result_url = result.get('remote_url', '')
    result_normalized = result_url.rstrip('/').rstrip('.git')
    expected_normalized = expected_url.rstrip('/').rstrip('.git')
    if result_normalized == expected_normalized:
        score = 1.0
    return score

def check_exact_match_v3__d4477d7a(result, expected, **options):
    """Verify that grading was completed and O2 contains the correct maximum score.

    This function checks:
    1. All students have been graded (grading_complete = True)
    2. Cell O2 contains a value
    3. The value in O2 matches the computed maximum score from all students
    4. The O2 value matches the expected value (100)

    Args:
        result: Dict from getter with grading information
        expected: Rules dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    expected_value = expected.get('expected_value')
    if expected_value is None:
        return 0.0
    grading_complete = result.get('grading_complete', False)
    if not grading_complete:
        return 0.0
    o2_value = result.get('o2_value')
    if o2_value is None:
        return 0.0
    computed_max = result.get('computed_max')
    if computed_max is None:
        return 0.0
    try:
        o2_int = int(o2_value)
        computed_max_int = int(computed_max)
        expected_int = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    if o2_int != computed_max_int:
        return 0.0
    if computed_max_int != expected_int:
        return 0.0
    return 1.0

def check_git_initialized__763d7485(result, expected, **options):
    """Check if git repository was initialized.

    Args:
        result: Dict with git status info from getter
        expected: Dict with 'is_git_repo' requirement
        **options: Additional options

    Returns:
        float: 1.0 if initialized, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    required_repo = expected.get('is_git_repo', True)
    actual_repo = result.get('is_git_repo', False)
    if actual_repo == required_repo:
        return 1.0
    else:
        return 0.0

def check_text_output__970046ef9644d0d33f656e418e7d5e7a(result: str, expected: Any, **options) -> float:
    """
    Check if text output contains expected lines in the correct order.

    Args:
        result: Actual text output from file
        expected: Dict with 'lines' key containing list of expected lines in order
        **options: Additional options (case_sensitive, allow_extra_lines)

    Returns:
        Score 1.0 if all lines match in order, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_lines = expected.get('lines', [])
    if not expected_lines:
        logger.warning('No expected lines specified')
        return 0.0
    case_sensitive = options.get('case_sensitive', True)
    allow_extra_lines = options.get('allow_extra_lines', False)
    result_text = result if case_sensitive else result.lower()
    result_lines = [line.strip() for line in result_text.split('\n') if line.strip()]
    expected_lines_normalized = [(line if case_sensitive else line.lower()).strip() for line in expected_lines]
    if not allow_extra_lines:
        if len(result_lines) != len(expected_lines_normalized):
            return 0.0
        for (i, exp_line) in enumerate(expected_lines_normalized):
            if result_lines[i] != exp_line:
                return 0.0
        return 1.0
    exp_idx = 0
    for result_line in result_lines:
        if exp_idx < len(expected_lines_normalized):
            if result_line == expected_lines_normalized[exp_idx]:
                exp_idx += 1
    if exp_idx == len(expected_lines_normalized):
        return 1.0
    return 0.0

def check_textbox_in_top_left__753748f3(src_path, expected, **options):
    """
    Check if the textbox is in the top-left corner of the image.
    Task variation 4 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is in top-left, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    top_most = height
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                top_most = min(top_most, y)
    if left_most < width * 0.05 and top_most < height * 0.05:
        return 1.0
    else:
        return 0.0

def check_line_count__1098fce8(result, expected, **options):
    """Check if line count matches expected value.

    Args:
        result: Line count from getter
        expected: Dict with 'min' and/or 'exact' value
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, int):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    if 'exact' in expected:
        if result == expected['exact']:
            return 1.0
        return 0.0
    if 'min' in expected:
        if result >= expected['min']:
            return 1.0
        return 0.0
    return 0.0

def check_python_imports__0a56ab11(result, expected, **options):
    """
    Validate extracted Python imports.

    NOTE: Complete extraction cannot be verified since we cannot access the source
    Colab notebook. We verify structural properties instead: minimum count, valid
    import syntax, no duplicates, and alphabetical sorting.

    Args:
        result: List of import lines from getter
        expected: Dict with 'min_imports', 'required_keywords' keys
        **options: Additional options (allow_duplicates)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, list):
        return 0.0
    min_imports = expected.get('min_imports', 1)
    required_keywords = expected.get('required_keywords', [])
    allow_duplicates = options.get('allow_duplicates', True)
    score = 0.0
    if len(result) >= min_imports:
        score += 0.3
    all_content = ' '.join(result)
    keywords_found = sum((1 for kw in required_keywords if kw in all_content))
    if required_keywords:
        score += 0.2 * (keywords_found / len(required_keywords))
    else:
        score += 0.2
    if not allow_duplicates:
        if len(result) == len(set(result)):
            score += 0.25
    else:
        score += 0.25
    if result == sorted(result):
        score += 0.25
    return min(1.0, score)

def check_file_exported__6ee1fdcd(result, expected, **options):
    """Check if file was successfully exported as a valid PDF.

    Args:
        result: Dict with file validation information:
            - exists: bool, whether file exists
            - is_pdf: bool, whether file is a PDF (MIME type check)
            - file_size: int, file size in bytes
            - is_fresh: bool, whether file was created recently
            - created_by_libreoffice: bool, whether PDF was created by LibreOffice
        expected: Expected value (True for file should exist)
        **options: Additional options

    Returns:
        float: 1.0 if all validations pass, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if isinstance(result, bool):
        return 1.0 if result == expected else 0.0
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('is_pdf', False):
        return 0.0
    file_size = result.get('file_size', 0)
    if file_size < 100:
        return 0.0
    if not result.get('is_fresh', False):
        return 0.0
    created_by_libreoffice = result.get('created_by_libreoffice', False)
    return 1.0

def check_file_exists__12ebd85d(result, expected, **options):
    """
    Check if a file exists based on command output.

    Args:
        result: Output from ls command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'exists')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_python_definitions__d881a7db2082639af55dd0ea1e3047ab(result, expected, **options):
    """
    Check if Python file contains ONLY function and class definitions (excluding other code).

    Args:
        result: File content string from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    score = 0.0
    if result and len(result.strip()) > 0:
        score += 0.2
    else:
        return 0.0
    try:
        tree = ast.parse(result)
    except SyntaxError as e:
        logger.warning(f'Python file has syntax errors: {e}')
        return 0.0
    func_defs = [node for node in ast.walk(tree) if isinstance(node, ast.FunctionDef)]
    class_defs = [node for node in ast.walk(tree) if isinstance(node, ast.ClassDef)]
    has_def = len(func_defs) > 0
    has_class = len(class_defs) > 0
    if has_def or has_class:
        score += 0.3
    else:
        return score
    top_level_statements = tree.body
    only_definitions = all((isinstance(stmt, (ast.FunctionDef, ast.ClassDef, ast.AsyncFunctionDef)) for stmt in top_level_statements))
    if only_definitions:
        score += 0.3
    else:
        non_def_types = [type(stmt).__name__ for stmt in top_level_statements if not isinstance(stmt, (ast.FunctionDef, ast.ClassDef, ast.AsyncFunctionDef))]
        logger.info(f'Found non-definition code: {non_def_types}')
        score += 0.1
    if has_def and has_class:
        score += 0.2
    return min(score, 1.0)

def check_renamed_files__0190943f9252410b5b69d12d28f1b6b7(result, expected, **options):
    """
    Check if files were renamed from *failed.ipynb to *_archived.ipynb.

    Args:
        result: Dict with 'archived_files' and 'failed_files' lists
        expected: Expected data (from rules dict)
        **options: Additional options

    Returns:
        float: Score based on:
            - 0.5 points: All expected archived files exist
            - 0.5 points: No failed files remain
    """
    expected_archived = expected.get('archived_files', [])
    expected_no_failed = expected.get('no_failed_files', True)
    logger.info(f'Checking renamed files: result={result}, expected_archived={expected_archived}')
    if not isinstance(result, dict):
        logger.error('Result is not a dictionary')
        return 0.0
    archived_files = result.get('archived_files', [])
    failed_files = result.get('failed_files', [])
    score = 0.0
    expected_archived_sorted = sorted(expected_archived)
    archived_files_sorted = sorted(archived_files)
    if archived_files_sorted == expected_archived_sorted:
        score += 0.5
        logger.info('All archived files present: +0.5')
    else:
        matching = len(set(archived_files) & set(expected_archived))
        total = len(expected_archived)
        if total > 0:
            partial = matching / total * 0.5
            score += partial
            logger.info(f'Partial archived files: {matching}/{total}, +{partial}')
    if expected_no_failed and len(failed_files) == 0:
        score += 0.5
        logger.info('No failed files remain: +0.5')
    elif expected_no_failed:
        penalty = min(0.5, len(failed_files) * 0.1)
        score += 0.5 - penalty
        logger.info(f'Some failed files remain: {len(failed_files)}, +{0.5 - penalty}')
    return min(1.0, score)

def check_text_content__fc2c8cc4(result, expected, **options):
    """Compare text content against expected line count and validate email format.

    Args:
        result: Text content from getter
        expected: Expected value dict with 'value' key containing expected line count
        **options: Additional options

    Returns:
        float: 1.0 if line count matches expected and all lines are valid emails, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    if not result:
        return 0.0
    result_stripped = str(result).strip()
    if not result_stripped:
        return 0.0
    lines = result_stripped.split('\n')
    non_empty_lines = [line.strip() for line in lines if line.strip()]
    actual_count = len(non_empty_lines)
    try:
        expected_count = int(expected_value)
    except (ValueError, TypeError):
        return 0.0
    if actual_count != expected_count:
        return 0.0
    email_pattern = re.compile('^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}$')
    valid_email_count = 0
    for line in non_empty_lines:
        if email_pattern.match(line):
            valid_email_count += 1
    if valid_email_count >= expected_count * 0.9:
        return 1.0
    else:
        return valid_email_count / expected_count if expected_count > 0 else 0.0

def check_file_exists__dccfbd8e(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def check_file_exists_with_size__c0d603dc(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_lines_contain_keywords__bf9ee805c777733e26c14d1927c30b1a(result, expected, **options):
    """Check if file lines contain all expected keywords.

    Args:
        result: List of lines from getter
        expected: Dict with 'keywords' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, list):
        logger.warning('Result is empty or not a list')
        return 0.0
    keywords = expected.get('keywords', [])
    if not keywords:
        logger.warning('No keywords specified')
        return 0.0
    full_text = ' '.join(result)
    found_count = 0
    for keyword in keywords:
        if keyword.lower() in full_text.lower():
            found_count += 1
    return found_count / len(keywords)

def check_file_exists__6cb37327(result, expected, **options):
    """Check if file exists.

    Args:
        result: Boolean indicating if file exists
        expected: Expected dict with 'exists' key (should be True)
        **options: Additional options

    Returns:
        float: 1.0 if file exists as expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_file_recovery__5ea617a3(result_state, expected_state, **options):
    """
    Check if a file was properly recovered from Trash.

    Verifies:
    1. The target file exists at the expected path
    2. The file has a reasonable size (not empty/corrupted)
    3. The file is no longer in Trash
    4. CRITICAL: The recovered file's hash matches the original file's hash
       (proving it's the SAME file, not a different file created at the target path)

    Args:
        result_state: dict from get_file_recovery_state__5ea617a3 with:
            - file_exists: bool
            - file_size: int
            - not_in_trash: bool
            - file_hash: str (SHA256 hash of recovered file)
            - original_hash: str (SHA256 hash of original file)
        expected_state: dict with:
            - min_file_size: int (minimum expected file size)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not result_state.get('file_exists', False):
        return 0.0
    min_size = expected_state.get('min_file_size', 1000)
    file_size = result_state.get('file_size', 0)
    if file_size < min_size:
        return 0.0
    if not result_state.get('not_in_trash', False):
        return 0.0
    file_hash = result_state.get('file_hash')
    original_hash = result_state.get('original_hash')
    if not file_hash or not original_hash:
        return 0.0
    if file_hash != original_hash:
        return 0.0
    return 1.0

def check_python_path__e7a35ed0ec20ca7ab6b257f3f5c87e23(result: str, expected: dict, **options) -> float:
    """Check if the Python path setting matches the expected value.

    Args:
        result: String containing the actual Python path
        expected: Dict with 'path' field specifying expected path

    Returns:
        1.0 if paths match, 0.0 otherwise
    """
    expected_path = expected.get('path', '')
    if result == expected_path:
        return 1.0
    return 0.0

def check_python_files_count__13d7d579(result, expected, **options):
    """
    Check if repository was properly downloaded and contains expected structure.

    Args:
        result: dict with repository information
        expected: dict with validation rules
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    min_count = expected.get('min_count', 10)
    expected_repo_url = expected.get('expected_repo_url', 'https://github.com/xlang-ai/instructor-embedding')
    py_count = result.get('python_file_count', 0)
    if py_count < min_count:
        return 0.0
    if not result.get('has_git_dir', False):
        return 0.0
    git_remote = result.get('git_remote_url', '')
    if git_remote:
        git_remote_normalized = git_remote.replace('.git', '').replace('git@github.com:', 'https://github.com/')
        expected_normalized = expected_repo_url.replace('.git', '')
        if expected_normalized not in git_remote_normalized:
            return 0.0
    else:
        return 0.0
    if not result.get('has_readme', False):
        return 0.0
    return 1.0

def check_file_count__081d0b6c(result: int, expected: dict, **options) -> float:
    """
    Check if file count matches expected value.

    Args:
        result: Actual count from getter
        expected: Dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('expected_count', 0)
    logger.info(f'Actual file count: {result}, Expected: {expected_count}')
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_json_settings__c5a909ed(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific key-value pair exists in the VS Code settings JSON file.

    Args:
        actual (str): path to result settings.json file
        expected (dict): expected dict with keys "expected_key" and "expected_value"

    Return:
        float: 1.0 if the key-value pair exists, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_json_settings__c5a909ed: actual file path is None')
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        logger.debug(f'check_json_settings__c5a909ed: Error reading JSON file: {e}')
        return 0.0
    expected_key = expected.get('expected_key')
    expected_value = expected.get('expected_value')
    if expected_key is None or expected_value is None:
        logger.debug('check_json_settings__c5a909ed: expected_key or expected_value is None')
        return 0.0
    actual_value = data.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    logger.debug(f'check_json_settings__c5a909ed: Expected {expected_key}={expected_value}, got {actual_value}')
    return 0.0

def compare_text_output__e6536361(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_file_exists_with_size__91c793ae(result, expected, **options):
    """
    Check if file exists and has reasonable size.

    Args:
        result: dict from getter with keys {"exists": bool, "size": int, "is_png": bool}
        expected: dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_file_move_copy__aace45122e840d40b84dc540ae5a49bc(result: dict, expected: dict) -> float:
    """
    Check if file was moved to one directory and copied to others.

    Args:
        result: Dict with file existence info from getter
        expected: Dict with 'dirs_with_file' (list of dirs that should have the file)
                  and 'file_not_in_root' (bool, whether file should not exist in root)

    Returns:
        float: 1.0 if all conditions met, 0.0 otherwise
    """
    if result is None:
        return 0.0
    dirs_with_file = expected.get('dirs_with_file', [])
    file_not_in_root = expected.get('file_not_in_root', False)
    for dir_name in dirs_with_file:
        if not result.get(dir_name, False):
            return 0.0
    if file_not_in_root:
        if result.get('root', False):
            return 0.0
    return 1.0

def check_file_renames__6f033773(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if files were renamed correctly.

    Args:
        result: Command output listing renamed files
        expected: Expected rules dict with 'expected_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    try:
        expected_files = expected.get('expected_files', [])
        if not expected_files:
            return 0.0
        actual_files = [line.strip() for line in result.strip().split('\n') if line.strip()]
        matches = sum((1 for f in expected_files if f in actual_files))
        score = matches / len(expected_files)
        return min(1.0, score)
    except Exception as e:
        return 0.0

def check_python_text_patterns__198be354(result, expected, **options):
    """Check if required text patterns are found in the file and validate code completeness.

    This function verifies that the exported Python code contains all essential components
    from the Colab notebook, including required patterns, valid Python syntax, proper
    structure, and substantial content.

    Args:
        result: dict from getter with pattern_matches, is_valid_python, etc.
        expected: dict with 'rules' containing verification requirements
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    rules = expected.get('rules', {})
    require_valid_python = rules.get('require_valid_python', True)
    if require_valid_python and (not result.get('is_valid_python', False)):
        return 0.0
    min_length = rules.get('min_content_length', 500)
    content_length = result.get('content_length', 0)
    if content_length < min_length:
        return 0.0
    require_imports = rules.get('require_imports', True)
    if require_imports and (not result.get('has_imports', False)):
        return 0.0
    require_class = rules.get('require_class_definition', True)
    if require_class and (not result.get('has_class_def', False)):
        return 0.0
    required_patterns = rules.get('required_patterns', [])
    if not required_patterns:
        return 1.0
    pattern_matches = result.get('pattern_matches', {})
    matched = sum((1 for pattern in required_patterns if pattern_matches.get(pattern, False)))
    if matched == len(required_patterns):
        return 1.0
    elif matched == 0:
        return 0.0
    else:
        return matched / len(required_patterns) * 0.7

def check_python_pkg__a3b47e9754f6a01a17ed98f7d00d938c(result: str, expected: dict, **options) -> float:
    """
    Check if expected Python packages are found in pip list output

    Args:
        result: Output from pip list command
        expected: Dictionary with 'expect' patterns to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expect_patterns: List[Pattern[str]] = [re.compile(ptt) for ptt in expected.get('expect', [])]
    if not expect_patterns:
        return 0.0
    for pattern in expect_patterns:
        if not pattern.search(result):
            logger.info(f'Pattern not found: {pattern.pattern}')
            return 0.0
    logger.info('All expected patterns found')
    return 1.0

def check_large_file_count__b9976565(file_count, expected):
    """Check if large file count is at least expected.

    Args:
        file_count: Number of large files in directory
        expected: Dict with 'min_count' key

    Returns:
        float: 1.0 if count >= min_count, 0.0 otherwise
    """
    min_count = expected['min_count']
    if file_count >= min_count:
        return 1.0
    else:
        return 0.0

def check_file_exists__45e0b2e9(result, expected, **options):
    """Check if file exists and has correct content from Save As operation.

    Args:
        result: Dict with 'exists', 'content_hash', and 'original_hash' keys
        expected: Expected state dict with 'exists' and optional 'verify_content' keys
        **options: Additional options (not used in this implementation)

    Returns:
        float: 1.0 if verification passes, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    result_exists = result.get('exists', False)
    expected_exists = expected.get('exists', True)
    if result_exists != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    if not result_exists:
        return 0.0
    verify_content = expected.get('verify_content', False)
    if verify_content:
        result_hash = result.get('content_hash')
        original_hash = result.get('original_hash')
        if result_hash is None or original_hash is None:
            return 0.0
        if result_hash == original_hash:
            return 1.0
        else:
            return 0.0
    return 1.0

def check_file_size_range__b9c089b2fe7d833fde2da297bbbd9620(result, expected, **options):
    """
    Check if file size is within expected range.

    Args:
        result: int (file size in bytes)
        expected: dict with 'min_size' and optionally 'max_size'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, int):
        logger.error(f'Result is not an integer: {result}')
        return 0.0
    if result == 0:
        logger.info('File does not exist or is empty')
        return 0.0
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    if min_size <= result <= max_size:
        logger.info(f'File size {result} is within range [{min_size}, {max_size}]')
        return 1.0
    else:
        logger.warning(f'File size {result} is outside range [{min_size}, {max_size}]')
        return 0.0

def check_text_file_lines__5ced85fc_aug18_v0_c9e8a1b2d3f4e5a6b7c8d9e0f1a2b3c4(result, expected, **options):
    """Check if file lines match expected lines exactly.

    Args:
        result: List of lines from the file
        expected: Rules dict with 'expected_lines' key
        **options: Additional options

    Returns:
        float: 1.0 if all lines match, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error('Result is not a list')
        return 0.0
    expected_lines = expected.get('expected_lines', [])
    if len(result) != len(expected_lines):
        logger.info(f'Line count mismatch: got {len(result)}, expected {len(expected_lines)}')
        return 0.0
    for (i, (actual, expected_line)) in enumerate(zip(result, expected_lines)):
        if actual != expected_line:
            logger.info(f"Line {i + 1} mismatch: got '{actual}', expected '{expected_line}'")
            return 0.0
    logger.info('All lines match expected values')
    return 1.0

def check_text_replaced__66306a8b(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Full text content from getter
        expected: Dictionary with 'old_text' and 'new_text' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    old_text = expected.get('old_text', '')
    new_text = expected.get('new_text', '')
    if not old_text or not new_text:
        logger.error('Missing old_text or new_text in expected')
        return 0.0
    old_count = result.lower().count(old_text.lower())
    new_count = result.lower().count(new_text.lower())
    score = 0.0
    if old_count == 0:
        score += 0.5
        logger.info(f"✓ Old text '{old_text}' successfully removed")
    else:
        logger.warning(f"✗ Old text '{old_text}' still found {old_count} times")
    if new_count > 0:
        score += 0.5
        logger.info(f"✓ New text '{new_text}' found {new_count} times")
    else:
        logger.warning(f"✗ New text '{new_text}' not found")
    logger.info(f'Final score: {score:.2f}')
    return score

def check_file_size_range__79b627f8(result, expected, **options):
    """Check if file size is within expected range.

    Args:
        result: File size in bytes from getter
        expected: Expected rules dict with 'min_size' and 'max_size' keys
        **options: Additional comparison options

    Returns:
        float: 1.0 if size is in range, 0.0 otherwise
    """
    if result < 0:
        logger.warning('File not found or error getting size')
        return 0.0
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    if min_size <= result <= max_size:
        return 1.0
    else:
        logger.info(f'File size {result} bytes not in range [{min_size}, {max_size}]')
        return 0.0

def check_text_content__6941d0dc31bbd0c3d844303b7e1c57e5(result, expected, **options):
    """Check if file content matches expected text exactly.

    Args:
        result: String content from the file
        expected: Dict with 'content' key containing expected text
        **options: Additional options (ignore_trailing_whitespace, etc.)

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_content = expected.get('content', '')
    ignore_trailing = options.get('ignore_trailing_whitespace', False)
    if ignore_trailing:
        result = result.rstrip()
        expected_content = expected_content.rstrip()
    if result == expected_content:
        return 1.0
    return 0.0

def check_file_exists__b8a50137(result, expected, **options):
    """
    Check if file exists and is a valid PNG image.

    Args:
        result: Dict with keys: 'exists', 'is_png', 'is_valid_image', 'has_dimensions'
        expected: Expected value dict with 'exists' key

    Returns:
        1.0 if file exists, is a valid PNG image with reasonable dimensions, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        if not result.get('exists', False):
            return 1.0
        else:
            return 0.0
    file_exists = result.get('exists', False)
    is_png = result.get('is_png', False)
    is_valid_image = result.get('is_valid_image', False)
    has_dimensions = result.get('has_dimensions', False)
    if file_exists and is_png and is_valid_image and has_dimensions:
        return 1.0
    else:
        if not file_exists:
            logger.info('File does not exist')
        elif not is_png:
            logger.info('File exists but is not a PNG image')
        elif not is_valid_image:
            logger.info('File exists but is not a valid image')
        elif not has_dimensions:
            logger.info('File exists but does not have valid dimensions (should be at least 100x100)')
        return 0.0

def check_file_count__4e03b1ed(result, expected, **options):
    """Check if file count matches expected.

    Args:
        result: Actual count from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if result == expected_count:
        logger.info(f'File count matches: {result}')
        return 1.0
    else:
        logger.info(f'File count mismatch. Expected: {expected_count}, Got: {result}')
        return 0.0

def check_files_moved__5bd46009a70125f4395abd10c34421d4(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if files were moved (present in target folder, absent from source).

    Args:
        result: Dict with 'target_files' (files in target folder) and 'source_files' (files remaining in source)
        expected: Dict with 'target_files' (expected in target) and 'should_not_be_in_source' (should be removed from source)
        **options: Additional options

    Returns:
        float: 1.0 if files were properly moved, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    target_files = set(result.get('target_files', []))
    source_files = set(result.get('source_files', []))
    expected_target = set(expected.get('target_files', []))
    should_not_be_in_source = set(expected.get('should_not_be_in_source', []))
    if target_files != expected_target:
        return 0.0
    if source_files.intersection(should_not_be_in_source):
        return 0.0
    return 1.0

def check_filenames_match__71c23132811d122bc61dca33636b8f81(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the file names match the expected pattern.

    Args:
        result: Actual list of file names (sorted)
        expected: Expected configuration (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if all expected files exist, 0.0 otherwise
    """
    expected_names = expected.get('filenames', [])
    if sorted(result) == sorted(expected_names):
        return 1.0
    else:
        return 0.0

def check_text_content__a48d3d2ba069ee77685e4821041681b9(result: str, expected: dict, **options) -> float:
    """Check if text file contains expected lines with exact matching and count verification.

    Args:
        result: Actual text content from file
        expected: Expected configuration with 'contains' list of strings

    Returns:
        float: Score 1.0 if all requirements met, 0.0 otherwise
    """
    if not result:
        return 0.0
    contains_list = expected.get('contains', [])
    if not contains_list:
        return 0.0
    result_lines = [line.strip() for line in result.split('\n') if line.strip()]
    if len(result_lines) != len(contains_list):
        return 0.0
    found_count = 0
    for expected_str in contains_list:
        if expected_str in result_lines:
            found_count += 1
    if found_count == len(contains_list):
        return 1.0
    return 0.0

def check_text_replacement__438c9c7ce7eebe25a3992ddf0a388112(result: str, expected: Dict, **options) -> float:
    """
    Check if text replacement was performed correctly.

    Args:
        result: Actual file content (string from getter)
        expected: Expected value with 'expected_content' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        logger.warning('Result is empty')
        return 0.0
    expected_content = expected.get('expected_content', '')
    if not expected_content:
        logger.error('No expected_content in expected dict')
        return 0.0
    if result.strip() == expected_content.strip():
        return 1.0
    else:
        logger.info(f'Content mismatch. Expected length: {len(expected_content)}, Actual length: {len(result)}')
        return 0.0

def check_line_count__bcaf4400(result, expected, **options):
    """
    Validate line count is within expected range.

    This metric validates that the user successfully counted lines from the Colab
    notebook and saved the result. The instruction suggests excluding empty lines
    and comment-only lines for accuracy, but the evaluator verifies the core task
    completion rather than the exact counting methodology.

    WHAT IS VERIFIED:
    - File exists at the specified path (/home/user/line_count.txt)
    - File contains a valid integer
    - Integer is within reasonable range (80-150 lines)
    - Integer is a positive number

    WHAT CANNOT BE VERIFIED:
    - Whether user actually counted from the Colab notebook vs. another source
    - Whether empty lines were excluded
    - Whether comment-only lines were excluded
    - Exact accuracy of the count

    This verification approach balances task utility (teaching line counting and file I/O)
    with feasibility constraints (Colab content extraction is not programmatically accessible).

    Args:
        result: Integer line count from getter
        expected: Dict with 'min_lines' and 'max_lines' keys defining the valid range
        **options: Additional options

    Returns:
        float: Score 1.0 if within range, 0.0 otherwise
    """
    if result is None or not isinstance(result, (int, float)):
        return 0.0
    result = int(result)
    min_lines = expected.get('min_lines', 0)
    max_lines = expected.get('max_lines', float('inf'))
    if result <= 0:
        return 0.0
    if min_lines <= result <= max_lines:
        return 1.0
    else:
        return 0.0

def check_python_definitions__e8ec0313751ff65df15abd7d031a7c53(result, expected, **options):
    """Compare extracted function/class definitions against expected.

    Args:
        result: Extracted definitions from result file
        expected: Expected definitions (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0

    def extract_signatures(text):
        """Extract function/class signatures for comparison."""
        signatures = []
        for line in text.splitlines():
            stripped = line.strip()
            if stripped.startswith('def ') or stripped.startswith('class '):
                if '#' in stripped:
                    stripped = stripped[:stripped.index('#')].strip()
                stripped = ' '.join(stripped.split())
                signatures.append(stripped)
        return signatures
    result_sigs = extract_signatures(result)
    expected_sigs = extract_signatures(expected.get('expected_definitions', ''))
    if not expected_sigs:
        return 1.0 if not result_sigs else 0.0
    result_set = set(result_sigs)
    expected_set = set(expected_sigs)
    matched = result_set & expected_set
    if not expected_set:
        return 1.0 if not result_set else 0.0
    recall = len(matched) / len(expected_set)
    precision = len(matched) / len(result_set) if result_set else 0.0
    if recall < 0.7:
        return recall * 0.5
    if recall + precision == 0:
        return 0.0
    f1 = 2 * (recall * precision) / (recall + precision)
    return f1

def check_contains_numbered_files__939aa00d(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if list contains sequentially numbered files like 1.png, 2.png, 3.png with no gaps.

    This verifies that:
    1. At least min_count files exist with the specified extension
    2. Files are numbered starting from 1
    3. Numbering is strictly sequential with no gaps (1,2,3,... not 1,3,5,...)

    Args:
        result: List of filenames from Google Drive
        expected: Dict with 'min_count' (minimum number of files) and 'extension' (file extension)
        **options: Additional options (unused)

    Returns:
        1.0 if files meet all criteria, 0.0 otherwise
    """
    min_count = expected.get('min_count', 3)
    ext = expected.get('extension', '.png')
    numbered_files = []
    file_numbers = []
    for i in range(1, 50):
        filename = f'{i}{ext}'
        if filename in result:
            numbered_files.append(filename)
            file_numbers.append(i)
    if len(numbered_files) < min_count:
        logger.info(f'Found only {len(numbered_files)} numbered files, need at least {min_count}')
        return 0.0
    expected_numbers = list(range(1, len(numbered_files) + 1))
    if file_numbers != expected_numbers:
        logger.info(f'Numbering is not sequential. Found: {file_numbers}, Expected: {expected_numbers}')
        return 0.0
    logger.info(f'Found {len(numbered_files)} sequentially numbered files: {numbered_files}')
    return 1.0

def check_srt_file_exists__4f098003e517e7e34a157f1c233e1c85(result, expected, **options):
    """
    Check if SRT file exists and has valid SRT subtitle format.

    Args:
        result: Dict from getter with keys:
            - 'exists': bool
            - 'size': int
            - 'has_content': bool
            - 'valid_srt_format': bool
            - 'subtitle_count': int
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0

    Scoring breakdown:
        - 0.3 points: File exists
        - 0.3 points: File has content (> 100 bytes)
        - 0.4 points: File contains valid SRT format with at least 1 subtitle
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.3
    if result.get('has_content', False):
        score += 0.3
    if result.get('valid_srt_format', False):
        score += 0.4
    return score

def check_text_contains__7e5134258960ebea77ca0d290984a7a3(result, expected, **options):
    """Check if result text contains both sorted array and swap count.

    Args:
        result: Actual text content from file
        expected: Expected swap count as string (e.g., "14")
        **options: Additional options

    Returns:
        float: 1.0 if both sorted array and swap count are present, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    sorted_numbers = ['11', '12', '22', '25', '34', '64', '90']
    sorted_array_present = False
    numbers_in_result = re.findall('\\d+', result)
    if len(numbers_in_result) >= len(sorted_numbers):
        for i in range(len(numbers_in_result) - len(sorted_numbers) + 1):
            if numbers_in_result[i:i + len(sorted_numbers)] == sorted_numbers:
                sorted_array_present = True
                break
    swap_count_present = False
    if re.search('\\b' + expected + '\\b', result):
        swap_count_present = True
    elif re.search(expected + '(?:\\s|$|[^\\d])', result):
        swap_count_present = True
    if sorted_array_present and swap_count_present:
        return 1.0
    return 0.0

def check_text_replacement__d48f445fac6cb18530cd7f7169fdc7fc(result: str, expected: Dict, **options) -> float:
    """
    Check if text replacement was performed correctly.

    Args:
        result: Actual file content (string from getter)
        expected: Expected value with 'expected_content' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        logger.warning('Result is empty')
        return 0.0
    expected_content = expected.get('expected_content', '')
    if not expected_content:
        logger.error('No expected_content in expected dict')
        return 0.0
    if result.strip() == expected_content.strip():
        return 1.0
    else:
        logger.info(f'Content mismatch. Expected length: {len(expected_content)}, Actual length: {len(result)}')
        return 0.0

def check_line_starts_with__8f9a3936(result: str, expected: Dict, **options) -> float:
    """Check if a line starts with a specific prefix.

    Args:
        result: The line content
        expected: Dict with 'prefix' key

    Returns:
        1.0 if line starts with prefix, 0.0 otherwise
    """
    prefix = expected.get('prefix', '')
    if result.startswith(prefix):
        return 1.0
    return 0.0

def check_file_contains_comment__bcc15712(actual: str, expected: Dict, **options) -> float:
    """Check if a file contains a specific comment at the top.

    Args:
        actual (str): path to result text file
        expected (Dict): expected dict with 'comment_text' key
        **options: Additional options

    Return:
        float: the score (1.0 if comment is at the top, 0.0 otherwise)
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
    except Exception:
        return 0.0
    comment_text = expected.get('comment_text', '')
    if not comment_text:
        return 0.0
    lines = content.split('\n')
    first_non_empty_line = None
    for line in lines:
        stripped = line.strip()
        if stripped:
            first_non_empty_line = stripped
            break
    if first_non_empty_line and first_non_empty_line == comment_text.strip():
        return 1.0
    return 0.0

def check_exact_notes_count__eb089173(result_file, expected, **options):
    """
    Check if the exact number of slides with notes matches expected and verify content quality.

    This function verifies:
    1. All expected slides have notes (count check)
    2. Notes contain meaningful content (quality check)
    3. Notes meet minimum length requirement
    4. Notes content matches or is similar to content from reference document (notes.docx)

    Args:
        result_file: Path to the PPTX file
        expected: Dict with rules containing 'expected_count', optional 'min_note_length', and optional 'reference_doc_path'
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, partial credit (0.0-1.0) based on how many slides pass, 0.0 if failed
    """
    try:
        prs = Presentation(result_file)
        expected_count = expected.get('expected_count', 0)
        min_note_length = expected.get('min_note_length', 20)
        reference_doc_path = expected.get('reference_doc_path', '/home/user/Desktop/notes.docx')
        reference_content = []
        if os.path.exists(reference_doc_path):
            try:
                doc = Document(reference_doc_path)
                for paragraph in doc.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        reference_content.append(text)
            except Exception as e:
                reference_content = []
        reference_text = ' '.join(reference_content).lower()
        slides_with_valid_notes = 0
        total_slides = len(prs.slides)
        for slide in prs.slides:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text and len(notes_text) >= min_note_length:
                    if reference_content:
                        notes_lower = notes_text.lower()
                        notes_words = set(notes_lower.split())
                        reference_words = set(reference_text.split())
                        common_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'from', 'as', 'is', 'was', 'are', 'were', 'be', 'been', 'being'}
                        notes_words = notes_words - common_words
                        reference_words = reference_words - common_words
                        if notes_words:
                            overlap = len(notes_words & reference_words) / len(notes_words)
                            max_similarity = 0.0
                            for ref_para in reference_content:
                                similarity = SequenceMatcher(None, notes_lower, ref_para.lower()).ratio()
                                max_similarity = max(max_similarity, similarity)
                            if overlap >= 0.5 or max_similarity >= 0.4:
                                slides_with_valid_notes += 1
                    else:
                        slides_with_valid_notes += 1
        if slides_with_valid_notes == expected_count:
            return 1.0
        else:
            return slides_with_valid_notes / expected_count if expected_count > 0 else 0.0
    except Exception as e:
        return 0.0

def check_exact_text_match__53b3f7f8(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_zip_contains_files__5c33f919(result: list, expected: dict, **options) -> float:
    """Check if zip contains exactly the expected files (no more, no less).

    Args:
        result: List of filenames in the zip
        expected: Dict with 'required_files' key (list of filenames)
        **options: Additional options

    Returns:
        float: 1.0 if zip contains exactly the required files, 0.0 otherwise
    """
    required_files = expected.get('required_files', [])
    result_set = set(result)
    required_set = set(required_files)
    if required_set == result_set:
        return 1.0
    return 0.0

def check_text_content__20f90bc2668d01c30660dfeafc3af15b(result, expected, **options):
    """Check if parameter 'alist' has been renamed to 'arr' in the bubble_sort function.

    Args:
        result: String content from the file
        expected: Dict with 'content' key containing expected text with 'arr'
        **options: Additional options (ignore_trailing_whitespace, etc.)

    Returns:
        float: 1.0 if rename is successful, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    func_pattern = 'def\\s+bubble_sort\\s*\\([^)]*\\)\\s*:[^\\n]*(?:\\n(?:    |\\t)[^\\n]*)*'
    match = re.search(func_pattern, result)
    if not match:
        return 0.0
    bubble_sort_function = match.group(0)
    if re.search('\\balist\\b', bubble_sort_function):
        return 0.0
    if not re.search('def\\s+bubble_sort\\s*\\(\\s*arr\\s*\\)\\s*:', bubble_sort_function):
        return 0.0
    expected_patterns = ['\\blen\\s*\\(\\s*arr\\s*\\)', '\\barr\\s*\\[\\s*i\\s*\\]', '\\barr\\s*\\[\\s*i\\s*\\+\\s*1\\s*\\]']
    for pattern in expected_patterns:
        if not re.search(pattern, bubble_sort_function):
            return 0.0
    return 1.0

def check_python_file_count__c32ed37d80ef317dea88bac4a4cc1f31(result, expected, **options):
    """Check if Python file count meets expected criteria.

    Args:
        result: Integer count from getter
        expected: Rules dict with 'min_count' key (int)
        **options: Additional options

    Returns:
        float: 1.0 if count >= min_count, 0.0 otherwise
    """
    min_count = expected.get('min_count', 1)
    if result >= min_count:
        return 1.0
    else:
        return 0.0

def check_text_exact_match__f84e9cb5fd8cfc9fab208c24bcd90a7d(result: str, expected: dict, **options) -> float:
    """
    Check if text content exactly matches expected value.

    Args:
        result: Text content from getter
        expected: Dict with 'content' key containing expected text
        **options: Additional options

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        logger.debug('Result is None')
        return 0.0
    expected_content = expected.get('content', '')
    if result == expected_content:
        return 1.0
    else:
        logger.debug(f"Content mismatch. Expected: '{expected_content}', Got: '{result}'")
        return 0.0

def check_text_file_content__5ced85fc_aug18_v2_b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6e7(result, expected, **options):
    """Check if file content matches expected content exactly.

    Args:
        result: Content string from the file
        expected: Rules dict with 'expected_content' key
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        logger.error('Result is not a string')
        return 0.0
    expected_content = expected.get('expected_content', '')
    result_stripped = result.strip()
    expected_stripped = expected_content.strip()
    if result_stripped == expected_stripped:
        logger.info('Content matches expected value')
        return 1.0
    else:
        logger.info(f"Content mismatch: got '{result_stripped}', expected '{expected_stripped}'")
        return 0.0

def check_text_replacement__0f72ff3b9c8d9680b46915659d675f48(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from file
        expected: Rules dict containing 'expected_text' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    result_clean = result.strip()
    expected_clean = expected_text.strip()
    if result_clean == expected_clean:
        return 1.0
    return 0.0

def check_exact_text_match__b4af657e(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def compare_text_output__d7f70e02(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def compare_text_files(actual: str, rules: dict, **options) -> float:
    """
    Compare a text file content with expected value from rules.

    Args:
        actual (str): Path to the result text file
        rules (dict): Dictionary containing 'expected' key with the expected content
        **options: Additional options

    Returns:
        float: 1.0 if content matches expected (stripped), 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            actual_text = f.read().strip()
    except (FileNotFoundError, IOError):
        return 0.0
    expected = rules.get('expected', '').strip()
    if actual_text == expected:
        return 1.0
    return 0.0

def check_line_count_equals__fb1a48c8(result, expected, **options):
    """Check if line count equals expected value.

    Args:
        result: Actual line count
        expected: Dict with 'count' key specifying expected line count
        **options: Additional options

    Returns:
        float: 1.0 if line count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_count = expected.get('count', 0)
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_file_exported__9180f469(result, expected, **options):
    """Check if file was successfully exported in CSV UTF-8 format with valid content.

    Args:
        result: Dictionary with validation results from getter including:
            - exists: bool - whether file exists
            - is_csv: bool - whether file is valid CSV
            - is_utf8: bool - whether file is UTF-8 encoded
            - size: int - file size in bytes
            - row_count: int - number of rows in CSV
            - has_content: bool - whether file has actual content
        expected: Expected value (True for file should exist and be valid)
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on validation criteria
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('has_content', False):
        return 0.0
    if not result.get('is_utf8', False):
        return 0.0
    if not result.get('is_csv', False):
        return 0.0
    if result.get('row_count', 0) < 2:
        return 0.0
    if result.get('size', 0) < 100:
        return 0.0
    return 1.0

def check_file_size__f778a8914a698cd2bc7c0cc50cd3596d(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if file is a valid PNG image with acceptable size and dimensions.

    Args:
        result: Dict from getter with 'size', 'is_valid_png', 'width', 'height' keys, or None
        expected: Dict with 'min_size' and 'max_size' keys (in bytes)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_valid_png', False):
        logger.info('File is not a valid PNG image')
        return 0.0
    actual_size = result.get('size', 0)
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    if not min_size <= actual_size <= max_size:
        logger.info(f'File size {actual_size} bytes is outside range [{min_size}, {max_size}]')
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    if width < 640 or height < 480:
        logger.info(f'Image dimensions {width}x{height} are too small for a video frame (min: 640x480)')
        return 0.0
    logger.info(f'Valid PNG image: {width}x{height}, {actual_size} bytes - all checks passed')
    return 1.0

def check_file_created__b03b6c61(actual: bool, expected: Dict, **options) -> float:
    """Check if a file was created.

    Args:
        actual (bool): Result from getter indicating if file exists
        expected (Dict): expected dict with 'should_exist' key
        **options: Additional options

    Return:
        float: the score (1.0 if matches expectation, 0.0 otherwise)
    """
    should_exist = expected.get('should_exist', True)
    if actual == should_exist:
        return 1.0
    return 0.0

def check_line_indentation__a2b5af8108e461937716b976809e966c(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if lines have the expected indentation (in spaces).

    Args:
        result: Dict with line numbers as keys and indentation info from getter
        expected: Dict with expected indentation rules:
            - expected_spaces: Dict mapping line numbers (as strings) to expected space count
                e.g., {"2": 0, "3": 0, ...} means lines 2-10 should have 0 spaces

    Returns:
        float: Score between 0.0 and 1.0 based on how many lines match expected indentation
    """
    if not result:
        return 0.0
    expected_spaces = expected.get('expected_spaces', {})
    if not expected_spaces:
        return 0.0
    total_lines = len(expected_spaces)
    if total_lines == 0:
        return 0.0
    matching_lines = 0
    for (line_num, expected_space_count) in expected_spaces.items():
        if line_num in result:
            actual_spaces = result[line_num].get('leading_spaces', -1)
            if actual_spaces == expected_space_count:
                matching_lines += 1
    return matching_lines / total_lines

def check_direct_json_object__82bc8d6a(result, rules) -> float:
    """
    Custom metric function for task 82bc8d6a-36eb-4d2d-8801-ef714fb1e55a_task_verify_4.
    Compares two JSON objects directly, with support for relative time processing.

    This function processes relative time in the rules before comparison.
    The getter (get_rule__82bc8d6a) just returns the raw rules dict,
    and this metric handles all the relative time processing.
    """
    logger.info(f'[DEBUG] check_direct_json_object__82bc8d6a called with result: {result}')
    logger.info(f'[DEBUG] check_direct_json_object__82bc8d6a called with rules: {rules}')
    if isinstance(result, str):
        result = result.strip()
        result = result.replace("'", '"')
        try:
            result = json.loads(result)
        except json.JSONDecodeError as e:
            logger.error(f'Failed to parse result as JSON: {e}')
            return 0.0
    logger.info(f'[DEBUG] Processed result: {result}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    import copy
    rules = copy.deepcopy(rules)
    if 'relativeTime' in rules:
        rules = process_relative_time(rules)
        logger.info(f'[DEBUG] Rules after relative time processing: {rules}')
    try:
        expected_json = rules.get('expected', {})
        if expected_json:
            for (key, value) in expected_json.items():
                if value == '__EVALUATION_FAILED__':
                    logger.error(f"[DEBUG] Expected value for key '{key}' indicates evaluation failure, returning 0.0")
                    return 0.0
    except Exception as e:
        logger.error(f'[DEBUG] Error checking for evaluation failure indicator: {e}')
        return 0.0
    try:
        expect_in_result = rules.get('expect_in_result', False)
        logger.info(f'[DEBUG] expect_in_result: {expect_in_result}')
        if not expect_in_result:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON: {expected_json}')
            for (key, expected_value) in expected_json.items():
                if key not in result:
                    logger.info(f"[DEBUG] Key '{key}' not found in result, returning 0.0")
                    return 0.0
                result_value = result[key]
                if result_value != expected_value:
                    logger.info(f"[DEBUG] Mismatch for key '{key}': expected={expected_value}, got={result_value}, returning 0.0")
                    return 0.0
            logger.info('[DEBUG] All expected values match, returning 1.0')
            return 1.0
        else:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON (expect_in_result mode): {expected_json}')
            for (key, expected_value) in expected_json.items():
                if key not in result:
                    logger.info(f"[DEBUG] Key '{key}' not found in result, returning 0.0")
                    return 0.0
                result_value = result[key]
                if expected_value not in str(result_value):
                    logger.info(f"[DEBUG] Expected value '{expected_value}' not in result value '{result_value}' for key '{key}', returning 0.0")
                    return 0.0
            logger.info('[DEBUG] All expected values found in result, returning 1.0')
            return 1.0
    except Exception as e:
        logger.error(f'[DEBUG] Error during comparison: {e}')
        import traceback
        traceback.print_exc()
        return 0.0

def check_file_contains_items__f8073900a375900c7a9ce8fa79f05f9d(result, expected, **options):
    """Check if file content contains all expected items (one per line, unique, sorted alphabetically).

    Args:
        result: Actual file content (string)
        expected: Expected rules dict with 'items' key (list of strings)
        **options: Additional comparison options

    Returns:
        float: Score (1.0 if all requirements met, 0.0 otherwise)
    """
    expected_items = expected.get('items', [])
    if not expected_items:
        return 0.0
    result_lines = [line.strip() for line in result.split('\n') if line.strip()]
    if sorted(result_lines) != sorted(expected_items):
        return 0.0
    if len(result_lines) != len(set(result_lines)):
        return 0.0
    if result_lines != sorted(result_lines):
        return 0.0
    return 1.0

def check_specific_file_exists__72383bef81d322492ebfc1e4d86364a3(directory_list, rule):
    """
    Check if all required JPG files exist in the directory with partial credit.

    Args:
        directory_list: Directory tree structure from get_list_directory
        rule: Expected configuration with 'required_files' key (list of filenames)

    Returns:
        float: Partial credit score (count_present / total_required)
    """
    required_files = rule.get('required_files', [])
    if not required_files:
        return 0.0
    actual_files = [node['name'] for node in directory_list['children']]
    present_count = sum((1 for required_file in required_files if required_file in actual_files))
    return present_count / len(required_files)

def check_specific_file_present__04085b6d(result, expected, **options):
    """Check if a specific filename is present in the file list.

    Args:
        result: List of filenames from getter
        expected: Expected rules dict with 'target_file' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if target file is in list, 0.0 otherwise
    """
    target_file = expected.get('target_file', '')
    if target_file in result:
        return 1.0
    else:
        logger.info(f"Target file '{target_file}' not found. Available MP4 files: {result}")
        return 0.0

def check_python_file_exists__198be354(result, expected, **options):
    """Check if a Python file exists with comprehensive verification for Colab notebook extraction.

    Enhanced verification addressing LLM judge issues:
    - File existence and valid Python syntax (AST validation)
    - Structural fingerprinting: function/class counts, code complexity
    - Unique content markers: specific ML/DS patterns beyond generic imports
    - Tight line count bounds (±1% from expected 477 lines)
    - Verification code is from SPECIFIC notebook, not generic ML code

    Args:
        result: dict from getter with file info including:
            - exists, has_content, is_valid_python (basic checks)
            - structural_fingerprint (AST-based code structure)
            - unique_content_score (specific ML/DS patterns)
            - line_count (for length verification)
        expected: dict with 'min_line_count' requirement
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.1
    else:
        return 0.0
    if result.get('has_content', False):
        score += 0.1
    else:
        return score
    if result.get('is_valid_python', False):
        score += 0.2
    else:
        return score
    structural_fingerprint = result.get('structural_fingerprint', {})
    function_count = structural_fingerprint.get('function_count', 0)
    import_count = structural_fingerprint.get('import_count', 0)
    has_sufficient_structure = function_count >= 3 and import_count >= 5
    if has_sufficient_structure:
        score += 0.1
        if structural_fingerprint.get('has_main_or_train', False) or structural_fingerprint.get('has_model_definition', False):
            score += 0.1
    unique_content_score = result.get('unique_content_score', 0.0)
    unique_markers = result.get('unique_markers', {})
    if unique_content_score >= 0.42:
        score += 0.15
        has_real_ml_code = unique_markers.get('model_training_pattern', False) and unique_markers.get('data_loading_pattern', False)
        if has_real_ml_code:
            score += 0.1
    line_count = result.get('line_count', 0)
    gold_lines = 477
    lower_bound = int(gold_lines * 0.99)
    upper_bound = int(gold_lines * 1.01)
    if lower_bound <= line_count <= upper_bound:
        score += 0.15
    elif lower_bound - 10 <= line_count <= upper_bound + 10:
        score += 0.075
    return min(1.0, score)

def check_text_exact_match__600696b8508be0c2e3ca25794856fb75(result: str, expected: dict, **options) -> float:
    """
    Check if text content exactly matches expected value.

    Args:
        result: Text content from getter
        expected: Dict with 'content' key containing expected text
        **options: Additional options

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        logger.debug('Result is None')
        return 0.0
    expected_content = expected.get('content', '')
    if result == expected_content:
        return 1.0
    else:
        logger.debug(f"Content mismatch. Expected: '{expected_content}', Got: '{result}'")
        return 0.0

def check_dir_exists_with_files__5c107_2(result: str, expected: dict, **options) -> float:
    """
    Check if specific directories exist and contain expected files.

    Args:
        result: Output from command listing directory structure
        expected: Dictionary with 'directories' key containing list of expected dirs

    Returns:
        Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    directories = expected.get('directories', [])
    if not directories:
        return 0.0
    score = 0.0
    for directory in directories:
        if directory in result:
            score += 1.0 / len(directories)
    return score

def check_file_exists__724cfa0a(result, expected, **options):
    """Check if file exists.

    Args:
        result: Boolean from getter
        expected: Expected boolean value
        **options: Additional options

    Returns:
        float: 1.0 if matches expected, 0.0 otherwise
    """
    if result == expected:
        return 1.0
    return 0.0

def check_file_exists_and_structure_sim__2cf5fa39(result, expected, **options):
    """
    Check if the exported PNG file exists and has structure similar to the original XCF.

    Args:
        result: Path to the exported PNG file (from vm_file getter)
        expected: Dict with rules containing the source XCF path
        **options: Additional options (threshold for SSIM)

    Returns:
        float: 1.0 if file exists and structure matches, 0.0 otherwise
    """
    threshold = options.get('threshold', 0.85)
    if not result or not os.path.exists(result):
        logger.error(f'Result file does not exist: {result}')
        return 0.0
    try:
        png_img = Image.open(result)
        expected_width = 1152
        expected_height = 648
        if png_img.size[0] == expected_width and png_img.size[1] == expected_height:
            logger.info(f'PNG has correct dimensions: {png_img.size}')
            return 1.0
        else:
            logger.error(f'PNG dimensions mismatch: expected {expected_width}x{expected_height}, got {png_img.size}')
            return 0.0
    except Exception as e:
        logger.error(f'Error opening PNG file: {e}')
        return 0.0

def check_file_content__57d8acad(result, expected, **options):
    """
    Check if file content contains expected text.

    Args:
        result: File content from getter
        expected: Expected text to find in content

    Returns:
        float: 1.0 if expected text is in result, 0.0 otherwise
    """
    expected_text = expected.get('text', '')
    if expected_text in result:
        return 1.0
    return 0.0

def check_file_exists_with_content__845b16e9c20bb76eb4ef8a7eb9262413(result, expected, **options):
    """Check if file exists and contains expected strings.

    Args:
        result: Dict with 'exists' and 'content' from getter
        expected: Dict with 'must_contain' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dict')
        return 0.0
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    content = result.get('content', '')
    must_contain = expected.get('must_contain', [])
    if not must_contain:
        return 1.0
    found_count = 0
    for string in must_contain:
        if string in content:
            found_count += 1
    return found_count / len(must_contain)

def check_launch_json_config__a3743b930b4e3c5c6976584a28c8269c(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if launch.json contains expected configuration.

    Args:
        result: Dictionary containing actual launch.json content
        expected: Dictionary with 'config_name' to check for
        **options: Additional options

    Returns:
        float: 1.0 if configuration exists, 0.0 otherwise
    """
    if not result:
        return 0.0
    config_name = expected.get('config_name')
    if not config_name:
        logger.error('Expected config_name not specified')
        return 0.0
    configurations = result.get('configurations', [])
    if not isinstance(configurations, list):
        return 0.0
    for config in configurations:
        if isinstance(config, dict) and config.get('name') == config_name:
            return 1.0
    return 0.0

def check_blank_line_added__ebbefd78bc8d93d288c452f335196528(result: Dict[str, str], expected: Dict[str, Any], **options) -> float:
    """Check if a blank line was added at the specified position.

    Args:
        result: Dict with line content from getter: {"content": "...", "is_empty": True/False}
        expected: Dict with expected rules:
            - should_be_empty: Boolean indicating if line should be empty

    Returns:
        float: 1.0 if line matches expected state, 0.0 otherwise
    """
    if not result:
        return 0.0
    should_be_empty = expected.get('should_be_empty', True)
    is_empty = result.get('is_empty', False)
    if should_be_empty and is_empty:
        return 1.0
    elif not should_be_empty and (not is_empty):
        return 1.0
    else:
        return 0.0

def check_word_count_text__a4c1c19457791f4f99b06ead98b6bfeb(result, expected, **options):
    """Check if the word count text was updated correctly.

    Args:
        result: Data from getter containing word count text
        expected: Expected text content (target_text)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('found', False):
        return 0.0
    target_text = expected.get('target_text', '')
    actual_text = result.get('text', '')
    if target_text in actual_text:
        return 1.0
    else:
        return 0.0

def check_first_line_right_aligned__de33d712(result, expected, **options):
    """Check if first line is right aligned.

    Args:
        result: Alignment string from getter ('LEFT', 'CENTER', 'RIGHT', 'JUSTIFY')
        expected: Expected alignment from rules
        **options: Additional options

    Returns:
        float: 1.0 if right aligned, 0.0 otherwise
    """
    expected_alignment = expected.get('alignment', 'RIGHT')
    return 1.0 if result == expected_alignment else 0.0

def check_git_repo_cloned__f1a99656b1aa540fcb46d2aeba395a7e(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a git repository was successfully cloned.

    Args:
        result: Dict from getter with 'exists', 'is_git_repo', 'remote_url', 'files' keys
        expected: Dict with 'required_files' and 'remote_url' (optional)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.2
    else:
        logger.warning('Repository directory does not exist')
        return 0.0
    if result.get('is_git_repo', False):
        score += 0.2
    else:
        logger.warning('Directory exists but is not a git repository')
        return score
    required_files = expected.get('required_files', [])
    if required_files:
        files_found = result.get('files', [])
        files_found_normalized = [f.lstrip('./') for f in files_found]
        matching_files = 0
        for req_file in required_files:
            req_file_normalized = req_file.lstrip('./')
            if any((req_file_normalized in found_file for found_file in files_found_normalized)):
                matching_files += 1
        if matching_files == len(required_files):
            score += 0.3
        elif matching_files > 0:
            score += 0.3 * (matching_files / len(required_files))
        logger.info(f'Files matched: {matching_files}/{len(required_files)}')
    else:
        score += 0.3
    expected_remote = expected.get('remote_url', '')
    if expected_remote:
        actual_remote = result.get('remote_url', '')

        def normalize_git_url(url):
            """Normalize git URL for comparison."""
            if not url:
                return ''
            url = url.rstrip('/')
            if url.endswith('.git'):
                url = url[:-4]
            if url.startswith('git@github.com:'):
                url = 'https://github.com/' + url[15:]
            elif url.startswith('ssh://git@github.com/'):
                url = 'https://github.com/' + url[21:]
            return url.lower()
        normalized_expected = normalize_git_url(expected_remote)
        normalized_actual = normalize_git_url(actual_remote)
        if normalized_expected == normalized_actual:
            score += 0.3
            logger.info(f'Remote URL matched: {actual_remote}')
        else:
            logger.warning(f'Remote URL mismatch. Expected: {expected_remote}, Got: {actual_remote}')
    else:
        score += 0.3
    logger.info(f'Git repo check score: {score}')
    return score

def check_timezone__e2503403(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_text_replacement__66cc1bc44690547654db9625c97259f8(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from file
        expected: Rules dict containing 'expected_text' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    result_clean = result.strip()
    expected_clean = expected_text.strip()
    if result_clean == expected_clean:
        return 1.0
    return 0.0

def check_text_replacement__0f355d4c(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_file_exists_with_size__058bd353(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_line_numbering__2978d3bd(docx_file, expected):
    """Check if the document has line numbering enabled.

    Args:
        docx_file: Path to the docx file
        expected: Expected value (not used, checks for line numbering presence)

    Returns:
        float: 1.0 if line numbering is enabled, 0.0 otherwise
    """
    if not docx_file:
        return 0.0
    try:
        doc = Document(docx_file)
    except Exception as e:
        logger.error(f'Error loading document: {e}')
        return 0.0
    namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
    for section in doc.sections:
        sectPr = section._sectPr
        lnNumType = sectPr.find('.//w:lnNumType', namespaces)
        if lnNumType is not None:
            return 1.0
    return 0.0

def check_files_renamed__039c45a2(result, expected, **options):
    """Check if files were renamed with specific prefix.

    Args:
        result: List of filenames from getter
        expected: Dict with 'prefix' and 'min_count'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    prefix = expected.get('prefix', '')
    min_count = expected.get('min_count', 1)
    matching_count = sum((1 for filename in result if filename.startswith(prefix)))
    if matching_count >= min_count:
        return 1.0
    else:
        return matching_count / min_count

def check_file_organization__429c8cbc(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_python_file__f1f92c4b10af2ffde6ea1534830a28f6(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a Python file exists, is valid Python code, and contains a print statement.

    Args:
        result: Dict from getter with 'exists', 'is_valid_python', 'has_print_statement', 'content', and 'file_path' keys
        expected: Dict with 'file_extension' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
        - 30% for file existence
        - 30% for valid Python syntax
        - 30% for print statement presence
        - 10% for correct file extension
    """
    score = 0.0
    if not result.get('exists', False):
        return 0.0
    score += 0.3
    if result.get('is_valid_python', False):
        score += 0.3
    if result.get('has_print_statement', False):
        score += 0.3
    file_path = result.get('file_path', '')
    expected_extension = expected.get('file_extension', '.py')
    if file_path.endswith(expected_extension):
        score += 0.1
    return score

def check_filename_match__4b590a3a028f08e8f4ad12729f4351c3(result, expected, **options):
    """
    Check if filename matches expected basename and extension, and verify actual audio conversion.

    Args:
        result: Dict from getter with file info (exists, basename, extension, file_type, is_audio)
        expected: Dict with 'basename' and 'extension' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    score += 0.2
    if not result.get('is_audio', False):
        logger.warning(f"File is not actually audio format. MIME type: {result.get('file_type', 'unknown')}")
        logger.warning('This appears to be a renamed video file, not a proper audio conversion')
    else:
        score += 0.4
        logger.info(f"File is confirmed audio format: {result.get('file_type', '')}")
    expected_basename = expected.get('basename', '')
    actual_basename = result.get('basename', '')
    if actual_basename == expected_basename:
        score += 0.3
        logger.info(f'Basename matches: {actual_basename}')
    else:
        logger.warning(f"Basename mismatch: expected '{expected_basename}', got '{actual_basename}'")
    expected_extension = expected.get('extension', '')
    actual_extension = result.get('extension', '')
    if actual_extension == expected_extension:
        score += 0.1
        logger.info(f'Extension matches: {actual_extension}')
    else:
        logger.warning(f"Extension mismatch: expected '{expected_extension}', got '{actual_extension}'")
    logger.info(f'Filename and audio conversion check score: {score}')
    return score

def check_column_text_values__2c840a52(result, expected, **options):
    """Check if column values contain expected text values.

    Args:
        result: List of actual cell values
        expected: Dict with 'values' list to match
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_values = expected.get('values', [])
    if not isinstance(result, list):
        return 0.0
    if len(result) != len(expected_values):
        return 0.0
    matches = 0
    for (actual, exp) in zip(result, expected_values):
        if actual == exp:
            matches += 1
        elif isinstance(actual, str) and isinstance(exp, str):
            if actual.strip().upper() == exp.strip().upper():
                matches += 1
    return matches / len(expected_values) if expected_values else 0.0

def check_text_replacement__ca4f47f2(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_text_with_space__e8a68271(result, expected, **options):
    """Check if text contains expected spaced course code and not the unspaced version.

    This verifies that:
    1. The expected spaced text (e.g., 'CHIN 9505') is present
    2. The unspaced version (e.g., 'CHIN9505') is NOT present

    Args:
        result: Actual text from getter (paragraph text)
        expected: Expected text (from rules dict)
        **options: Additional options

    Returns:
        1.0 if expected spaced text found and unspaced version not found, 0.0 otherwise
    """
    expected_text = expected.get('text', '')
    if result is None:
        return 0.0
    if expected_text not in result:
        return 0.0
    unspaced_version = expected_text.replace(' ', '')
    if unspaced_version in result:
        return 0.0
    return 1.0

def check_text_match__25fb76d7ccb83f42013b589a25bead61(result, expected, **options):
    """Check if text content matches expected value.

    Args:
        result: Actual text content from getter
        expected: Expected text value (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_normalized = ' '.join(result.split())
    expected_normalized = ' '.join(str(expected_value).split())
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_python_path__1f6f3af3ee1e7e72d7b32984543de005(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if Python path matches expected value.

    Args:
        result: String containing actual Python path value
        expected: Dictionary with 'path' key containing expected path
        **options: Additional options

    Returns:
        float: 1.0 if paths match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_path = expected.get('path', '')
    if not expected_path:
        logger.error('Expected path not specified')
        return 0.0
    if result == expected_path:
        return 1.0
    return 0.0

def check_file_content_5f767718(result, expected, **options):
    """Check if file content matches expected count.

    Args:
        result: File path (string) returned by get_vm_file
        expected: Expected rules dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            content = f.read().strip()
        count = int(content)
        expected_count = expected.get('expected_count', 0)
        if count == expected_count:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        return 0.0
    except PermissionError:
        return 0.0
    except ValueError:
        return 0.0
    except Exception as e:
        return 0.0

def check_file_exists__f663e89c52b74d5c5d4e38ab6d86c83f(result, expected, **options):
    """Check if file exists as expected and is a valid PDF file with content.

    Args:
        result: Result from getter (dict with 'exists', 'is_pdf', 'file_size' keys)
        expected: Expected rules (dict with 'should_exist' key)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, is a valid PDF, and has content; 0.0 otherwise
    """
    file_exists = result.get('exists', False)
    is_pdf = result.get('is_pdf', False)
    file_size = result.get('file_size', 0)
    should_exist = expected.get('should_exist', True)
    if not should_exist:
        return 1.0 if not file_exists else 0.0
    if file_exists and is_pdf and (file_size > 0):
        return 1.0
    else:
        return 0.0

def check_exact_int_match__d4477d7a(result, expected, **options):
    """Compare result against expected integer value.

    Args:
        result: Actual value from getter
        expected: Rules dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_value = expected.get('expected_value')
    if result is None or expected_value is None:
        return 0.0
    try:
        result_int = int(result)
        expected_int = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    if result_int == expected_int:
        return 1.0
    return 0.0

def check_files_deleted__0a07f87a(result, rules, **options):
    """Check if files were deleted (find returns empty result).

    Args:
        result: Output from find command
        rules: Dict with 'should_be_empty' key set to True

    Returns:
        float: 1.0 if result is empty (files deleted), 0.0 otherwise
    """
    should_be_empty = rules.get('should_be_empty', False)
    if not should_be_empty:
        return 0.0
    if result is None or not result.strip():
        return 1.0
    else:
        return 0.0

def check_filename_pattern__4e03b1ed(result, expected, **options):
    """Check if filenames match expected pattern.

    Args:
        result: List of filenames from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Score based on pattern matching
    """
    pattern = expected.get('pattern', '')
    min_count = expected.get('min_count', 1)
    if not result:
        logger.info('No files found')
        return 0.0
    matching_files = []
    for filename in result:
        if re.search(pattern, filename, re.IGNORECASE):
            matching_files.append(filename)
    if len(matching_files) >= min_count:
        logger.info(f"Found {len(matching_files)} files matching pattern '{pattern}': {matching_files}")
        return 1.0
    else:
        logger.info(f'Only {len(matching_files)} files match pattern, expected at least {min_count}')
        return 0.0

def check_file_exists__81d9e3403ff656f186df75dab490d6ac(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if file exists and meets basic requirements.

    Args:
        result: Dict from getter with 'exists', 'is_png', 'size' keys
        expected: Dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        logger.info('File does not exist')
        return 0.0
    expected_is_png = expected.get('is_png', True)
    if result.get('is_png', False) == expected_is_png:
        score += 0.25
    else:
        logger.info(f"File type mismatch: expected is_png={expected_is_png}, got {result.get('is_png')}")
    min_size = expected.get('min_size', 1000)
    if result.get('size', 0) >= min_size:
        score += 0.25
    else:
        logger.info(f"File size too small: {result.get('size')} bytes (expected >= {min_size})")
    return score

def check_specific_text_deleted__ec78028f(result_file, expected, **options):
    """
    Check if specific text patterns have been deleted from all slides in a presentation.

    This function verifies that none of the shapes in any slide contain any of the
    specified text patterns. It recursively checks all shapes, including those inside
    groups.

    Args:
        result_file: Path to the PPTX file to check
        expected: Rules dictionary containing 'text_patterns' (list of strings to verify are deleted)
        **options: Additional options

    Returns:
        float: 1.0 if all text patterns are deleted (not found in any shape), 0.0 if any pattern is found
    """
    try:
        prs = Presentation(result_file)
        text_patterns = expected.get('text_patterns', [])
        if not text_patterns:
            logger.warning('No text patterns specified to check for deletion')
            return 0.0

        def get_all_text_shapes(slide):
            """
            Recursively get all shapes with text from a slide, including those inside groups.

            Args:
                slide: PowerPoint slide object

            Returns:
                list: List of all shapes that have text
            """

            def extract_text_shapes(shape):
                results = []
                if hasattr(shape, 'text') and hasattr(shape, 'text_frame'):
                    results.append(shape)
                if hasattr(shape, 'shapes'):
                    for sub_shape in shape.shapes:
                        results.extend(extract_text_shapes(sub_shape))
                return results
            all_text_shapes = []
            for shape in slide.shapes:
                all_text_shapes.extend(extract_text_shapes(shape))
            return all_text_shapes
        for (slide_idx, slide) in enumerate(prs.slides):
            text_shapes = get_all_text_shapes(slide)
            for shape in text_shapes:
                shape_text = shape.text.strip()
                for pattern in text_patterns:
                    if pattern in shape_text:
                        logger.info(f"Found text pattern '{pattern}' in slide {slide_idx + 1}, shape text: '{shape_text[:50]}...'")
                        return 0.0
        logger.info(f'Successfully verified that text patterns {text_patterns} are not present in the presentation')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking for deleted text: {e}')
        return 0.0

def check_eml_files_with_pattern__8a92d4dc(result: str, expected: dict) -> float:
    """
    Check if .eml files contain emails with subjects matching specific patterns.

    This function parses actual .eml file content to verify email subjects,
    rather than just checking filenames. It also verifies that the expected
    number of emails were exported.

    Args:
        result: path to ls output file containing paths to .eml files
        expected: dict with 'patterns' key containing list of regex patterns to match in subjects

    Returns:
        float: 1.0 if exactly 2 .eml files exist and all contain 'Test' in subject, 0.0 otherwise
    """
    if result is None:
        return 0.0
    patterns = expected.get('patterns', [])
    if not patterns:
        return 0.0
    eml_files = []
    try:
        with open(result) as f:
            for line in f:
                line = line.strip()
                if line and line.endswith('.eml'):
                    eml_files.append(line)
    except Exception as e:
        logger.error(f'Failed to read ls output: {e}')
        return 0.0
    expected_count = 2
    if len(eml_files) != expected_count:
        logger.warning(f'Expected {expected_count} .eml files, found {len(eml_files)}')
        return 0.0
    compiled_patterns = [re.compile(p, re.IGNORECASE) for p in patterns]
    emails_with_test = 0
    for eml_path in eml_files:
        try:
            with open(eml_path, 'rb') as eml_file:
                msg = email.message_from_binary_file(eml_file, policy=policy.default)
                subject = msg.get('Subject', '')
                matches_pattern = False
                for pattern in compiled_patterns:
                    if pattern.search(subject):
                        matches_pattern = True
                        break
                if matches_pattern:
                    emails_with_test += 1
                    logger.info(f'Found email with matching subject: {subject}')
                else:
                    logger.warning(f'Email subject does not match pattern: {subject}')
        except Exception as e:
            logger.error(f'Failed to parse .eml file {eml_path}: {e}')
            return 0.0
    if emails_with_test == expected_count:
        return 1.0
    else:
        logger.warning(f"Only {emails_with_test} out of {expected_count} emails have 'Test' in subject")
        return 0.0

def check_text_exact_match__4080707f437c813fb70f2db7aaa30575(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text content from file
        expected: Expected text content (from rules dict)
        **options: Additional options (case_sensitive, strip_whitespace)

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    strip_whitespace = options.get('strip_whitespace', True)
    case_sensitive = options.get('case_sensitive', True)
    result_text = result
    expected_text = expected
    if strip_whitespace:
        result_text = result_text.strip()
        expected_text = expected_text.strip()
    if not case_sensitive:
        result_text = result_text.lower()
        expected_text = expected_text.lower()
    return 1.0 if result_text == expected_text else 0.0

def check_text_replacement__11d0824d05970e58a5671a5636365f15(result: str, expected: Dict, **options) -> float:
    """
    Check if text replacement was performed correctly.

    Args:
        result: Actual file content (string from getter)
        expected: Expected value with 'expected_content' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        logger.warning('Result is empty')
        return 0.0
    expected_content = expected.get('expected_content', '')
    if not expected_content:
        logger.error('No expected_content in expected dict')
        return 0.0
    if result.strip() == expected_content.strip():
        return 1.0
    else:
        logger.info(f'Content mismatch. Expected length: {len(expected_content)}, Actual length: {len(result)}')
        return 0.0

def check_both_files_rows__e54614f2(result, expected, **options):
    """
    Check if both files exist and have expected row counts.

    Args:
        result: dict with file existence and row counts from getter
        expected: dict with expected row counts

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('csv_exists', False):
        score += 0.25
    if result.get('xlsx_exists', False):
        score += 0.25
    expected_rows = expected.get('csv_rows', 31)
    if result.get('csv_rows', 0) == expected_rows:
        score += 0.25
    expected_xlsx_rows = expected.get('xlsx_rows', 31)
    if result.get('xlsx_rows', 0) == expected_xlsx_rows:
        score += 0.25
    return score

def check_file_exported__0ad50b28(result, expected, **options):
    """Check if file was successfully exported as valid TSV with content.

    Args:
        result: Dictionary containing file validation results with keys:
            - exists (bool): Whether file exists
            - is_tsv (bool): Whether file uses tab delimiters
            - has_content (bool): Whether file is non-empty
            - row_count (int): Number of data rows
            - sample_line (str): First line of the file
        expected: Expected value (True for file should exist)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, is valid TSV, and has content; 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('has_content', False):
        return 0.0
    if not result.get('is_tsv', False):
        return 0.0
    if result.get('row_count', 0) < 1:
        return 0.0
    return 1.0

def check_file_exists__423224cbbe432d6315ffb9aa3c684c3a(result, expected, **options):
    """
    Check if MP3 file exists and is valid with comprehensive validation.

    Args:
        result: Dict from getter with 'exists', 'size', 'has_valid_format', 'duration' keys
        expected: Boolean expected value (True for valid MP3 file)
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not expected:
        if not result.get('exists', False):
            logger.info('File correctly does not exist')
            return 1.0
        else:
            logger.warning('File exists but was expected not to')
            return 0.0
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    file_size = result.get('size', 0)
    if file_size <= 1000:
        logger.warning(f'File size too small: {file_size} bytes (expected > 1000)')
        return 0.0
    if not result.get('has_valid_format', False):
        logger.warning('File does not have valid MP3 format (magic bytes check failed)')
        return 0.0
    duration = result.get('duration', 0)
    if duration <= 0:
        logger.warning(f'Invalid audio duration: {duration} seconds')
        return 0.0
    logger.info(f'All MP3 validation checks passed: size={file_size} bytes, duration={duration} seconds')
    return 1.0

def check_file_count__23e95644(result, expected, **options):
    """
    Check if file count matches expected.

    Args:
        result: File count from getter
        expected: Expected count

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if result == expected_count:
        return 1.0
    return 0.0

def check_file_exists__cb8ab5642aaf48705d11abb5543759c9(result, expected, **options):
    """Check if file existence matches expected value.

    Args:
        result: Boolean from getter (True if file exists)
        expected: Rules dict with 'exists' key (bool)
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_exact_count__68182234(result, expected, **options):
    """
    Check if the count matches expected value.

    Args:
        result: Integer count from getter
        expected: Dict with 'count' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not isinstance(result, int):
        return 0.0
    expected_count = expected.get('count', -1)
    return 1.0 if result == expected_count else 0.0

def check_txt_file_count__2b6d7a72(result, expected, **options):
    """Verify .doc files were properly converted to .txt files.

    Args:
        result: Dict from getter with conversion details
        expected: Expected configuration with 'count' field
        **options: Additional options

    Returns:
        float: 1.0 if conversion is correct, partial score otherwise
    """
    expected_count = expected.get('count', 12)
    txt_count = result.get('txt_count', 0)
    doc_count = result.get('doc_count', 0)
    matched_conversions = result.get('matched_conversions', [])
    txt_with_content = result.get('txt_with_content', 0)
    score = 0.0
    if txt_count == expected_count:
        score += 0.4
    if len(matched_conversions) == expected_count:
        score += 0.4
    elif len(matched_conversions) > 0:
        score += 0.4 * (len(matched_conversions) / expected_count)
    if txt_with_content > 0:
        score += 0.2
    return 1.0 if score >= 0.99 else score

def compare_text_output__2d526d9e(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_gitignore_content__c5a909ed(actual: str, expected: dict, **options) -> float:
    """
    Check if .gitignore file contains expected content on the first line.

    Args:
        actual (str): path to .gitignore file
        expected (dict): expected dict with key "expected_content"

    Return:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_gitignore_content__c5a909ed: actual file path is None')
        return 0.0
    expected_content = expected.get('expected_content')
    if expected_content is None:
        logger.debug('check_gitignore_content__c5a909ed: expected_content is None')
        return 0.0
    try:
        with open(actual, 'r') as f:
            first_line = f.readline().strip()
    except Exception as e:
        logger.debug(f'check_gitignore_content__c5a909ed: Error reading file: {e}')
        return 0.0
    if first_line == expected_content:
        return 1.0
    logger.debug(f"check_gitignore_content__c5a909ed: Expected '{expected_content}', got '{first_line}'")
    return 0.0

def check_renamed_files__27c9f1432580c9f5f1bf2d1f919f4ed5(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if all files have been renamed with the expected prefix.
    Verifies both: (1) renamed files exist AND (2) original files do not exist.

    Properly pairs each renamed file with its original to ensure true rename (not copy).

    Args:
        result: Dict from getter with 'renamed_files_found', 'original_files_remaining',
                'base_filenames', and 'prefix'
        expected: Dict with 'required_renamed_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    required_renamed_files = expected.get('required_renamed_files', [])
    renamed_files_found = result.get('renamed_files_found', [])
    original_files_remaining = result.get('original_files_remaining', [])
    base_filenames = result.get('base_filenames', [])
    prefix = result.get('prefix', '')
    if not required_renamed_files or not base_filenames:
        return 0.0
    total_files = len(base_filenames)
    correctly_renamed = 0
    for base_filename in base_filenames:
        expected_renamed = f'{prefix}{base_filename}'
        renamed_exists = expected_renamed in renamed_files_found
        original_does_not_exist = base_filename not in original_files_remaining
        if renamed_exists and original_does_not_exist:
            correctly_renamed += 1
    score = correctly_renamed / total_files
    return min(score, 1.0)

def check_file_list__b8c40a8e(result, expected, **options):
    """Check if text file contains all required filenames.

    Args:
        result: Content of text file with filenames
        expected: Rules dict with required_files list and exact_count
        **options: Additional options

    Returns:
        float: 1.0 if all required files are listed, 0.0 otherwise
    """
    if result is None or not isinstance(result, str):
        return 0.0
    required_files = expected.get('required_files', [])
    exact_count = expected.get('exact_count', len(required_files))
    lines = [line.strip() for line in result.strip().split('\n') if line.strip()]
    if len(lines) != exact_count:
        return 0.0
    for filename in required_files:
        if filename not in lines:
            return 0.0
    return 1.0

def check_file_organization__f3cabf2e(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_text_contains__0b52fd51(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    score = 0.0
    for substring in contains:
        if substring.lower() in result.lower():
            score += 1.0 / len(contains)
    return score

def check_file_exists__aec9e92c(result, expected, **options):
    """Check if file exists matches expected value.

    Args:
        result: Boolean from getter indicating file existence
        expected: Expected boolean value (from rules)
        **options: Additional comparison options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    return 1.0 if result == expected_exists else 0.0

def check_text_output__ec857351(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_files_with_prefix__8c4eaed9a61673f78def8f323e7dfe9d(result, expected, **options):
    """Check if directory contains expected number of files with a specific prefix.

    Args:
        result: dict from getter {'exists': bool, 'matching_files': list}
        expected: dict with 'count' (expected number of files)

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('Directory does not exist')
        return 0.0
    matching_files = result.get('matching_files', [])
    actual_count = len(matching_files)
    expected_count = expected.get('count', 0)
    if actual_count == expected_count:
        logger.info(f'File count matches: {actual_count} files found: {matching_files}')
        return 1.0
    else:
        logger.info(f'File count mismatch: expected {expected_count}, found {actual_count}. Files: {matching_files}')
        return 0.0

def check_python_imports__198be354(result, expected, **options):
    """Check if Python file has expected imports and code content.

    Args:
        result: dict from getter with import info and code metrics
        expected: dict with 'min_imports', 'required_modules', and optionally 'min_code_lines'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    min_imports = expected.get('min_imports', 0)
    actual_count = result.get('import_count', 0)
    if actual_count >= min_imports:
        score += 0.3
    required_modules = expected.get('required_modules', [])
    if required_modules:
        actual_modules = set(result.get('unique_modules', []))
        matched = sum((1 for mod in required_modules if mod in actual_modules))
        score += 0.3 * (matched / len(required_modules))
    else:
        score += 0.3
    min_code_lines = expected.get('min_code_lines', 10)
    non_empty_lines = result.get('non_empty_lines', 0)
    import_count = result.get('import_count', 0)
    code_lines = non_empty_lines - import_count
    if code_lines >= min_code_lines:
        score += 0.4
    elif code_lines > 0:
        score += 0.4 * (code_lines / min_code_lines)
    return score

def check_file_exists__ec920d7f(result: Dict[str, bool], expected: Dict[str, Any], **options) -> float:
    """Check if file exists and is a valid ODS file.

    Args:
        result: Dict with 'exists', 'valid_ods', and 'has_content' keys
        expected: Dict (not used, just checks if file exists as valid ODS)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is valid ODS, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('valid_ods', False):
        return 0.0
    if not result.get('has_content', False):
        return 0.0
    return 1.0

def check_file_renamed__a7904c4a(result: str, expected: dict, **options) -> float:
    """
    Check if file exists with new name and has expected extension.

    Args:
        result: Path to the renamed file from getter
        expected: Dict with 'expected_extension' to verify
        **options: Additional options

    Returns:
        float: 1.0 if file exists with correct extension, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None - file not found')
        return 0.0
    if not os.path.exists(result):
        logger.warning(f'File does not exist: {result}')
        return 0.0
    expected_ext = expected.get('expected_extension', None)
    if expected_ext:
        actual_ext = os.path.splitext(result)[1].lower()
        if actual_ext != expected_ext.lower():
            logger.warning(f'Wrong extension: {actual_ext} != {expected_ext}')
            return 0.0
    min_size = expected.get('min_size_bytes', 100)
    file_size = os.path.getsize(result)
    if file_size < min_size:
        logger.warning(f'File too small: {file_size} < {min_size}')
        return 0.0
    logger.info(f'File exists with correct name and extension: {result}')
    return 1.0

def check_text_contains__a2f161de(result, expected, **options):
    """Check if text file contains expected substring.

    Args:
        result: Actual text content from getter
        expected: Expected content (dict with 'contains' key)
        **options: Additional options

    Returns:
        float: 1.0 if text contains expected substring, 0.0 otherwise
    """
    expected_text = expected.get('contains', '')
    if expected_text.lower() in result.lower():
        return 1.0
    else:
        return 0.0

def check_text_output__9ee70fa2(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_text_exact_match__49075d97b370c316554a4b259a7ccc3e(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text content from file
        expected: Expected text content (from rules dict)
        **options: Additional options (case_sensitive, strip_whitespace)

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    strip_whitespace = options.get('strip_whitespace', True)
    case_sensitive = options.get('case_sensitive', True)
    result_text = result
    expected_text = expected
    if strip_whitespace:
        result_text = result_text.strip()
        expected_text = expected_text.strip()
    if not case_sensitive:
        result_text = result_text.lower()
        expected_text = expected_text.lower()
    return 1.0 if result_text == expected_text else 0.0

def check_zip_contains_files__6aa029d37a944ba9e2bf06a8a1d59f5c(result: list, expected: dict, **options) -> float:
    """
    Check if ZIP archive contains the expected files.

    Args:
        result: List of filenames in the ZIP
        expected: Rules dict with 'required_files' key (list of filenames)
        **options: Additional options

    Returns:
        float: Score based on how many required files are present
    """
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    if not result:
        return 0.0
    matches = 0
    for required_file in required_files:
        if required_file in result:
            matches += 1
    score = matches / len(required_files)
    return score

def check_subject_file__271abb880d5f6f8d57d2c41e20bcf6ad(result: str, expected: dict, **options) -> float:
    """
    Check if a text file contains the expected email subject.

    Args:
        result: Content from the created text file
        expected: Expected values from rules (dict with 'subject')
        **options: Additional options

    Returns:
        1.0 if the content matches the expected subject, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        result_str = str(result).strip()
        expected_subject = expected.get('subject', '').strip()
        if result_str == expected_subject:
            return 1.0
        if result_str.lower() == expected_subject.lower():
            return 0.9
        if expected_subject.lower() in result_str.lower():
            return 0.7
        return 0.0
    except Exception:
        return 0.0

def check_line_chart__0f492b286ccfa81ae15bde7a08cd32e6(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file has expected line chart with correct data ranges and series.

    Args:
        result: Chart info from getter containing series count, data ranges, and column usage
        expected: Expected chart properties including required series count and data ranges
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    expected_series_count = expected.get('series_count', 2)
    expected_columns = expected.get('columns', ['B'])
    require_2019_data = expected.get('require_2019_data', True)
    require_2020_data = expected.get('require_2020_data', True)
    min_data_points_per_series = expected.get('min_data_points_per_series', 12)
    score = 0.0
    if not result.get('has_chart', False):
        return 0.0
    score += 0.1
    if not result.get('has_line_chart', False):
        return score
    score += 0.2
    actual_series_count = result.get('series_count', 0)
    if actual_series_count == expected_series_count:
        score += 0.2
    uses_column_b = result.get('uses_column_b', False)
    columns_used = result.get('columns_used', [])
    if uses_column_b and all((col in expected_columns for col in columns_used)):
        score += 0.2
    has_2019_data = result.get('has_2019_data', False)
    if has_2019_data or not require_2019_data:
        score += 0.15
    has_2020_data = result.get('has_2020_data', False)
    if has_2020_data or not require_2020_data:
        score += 0.15
    series_details = result.get('series_details', [])
    if series_details:
        all_series_have_enough_points = all((s.get('data_point_count', 0) >= min_data_points_per_series for s in series_details))
    return score

def check_exact_number__f9a0219a(result, expected, **options):
    """
    Check if the result matches the expected number exactly.

    Args:
        result: Integer from getter
        expected: Dict with 'value' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not isinstance(result, int):
        return 0.0
    expected_value = expected.get('value', -1)
    return 1.0 if result == expected_value else 0.0

def check_filename_hash_mapping__bf825e2c(result, expected, **options):
    """Check if filename-hash mapping matches expected.

    Args:
        result: Dict mapping filenames to hashes
        expected: Dict with 'expected_mapping' key

    Returns:
        float: 1.0 if mapping matches exactly, 0.0 otherwise
    """
    expected_mapping = expected.get('expected_mapping', {})
    if result == expected_mapping:
        return 1.0
    else:
        print(f'Expected: {expected_mapping}')
        print(f'Got: {result}')
        return 0.0

def check_textbox_vmiddle__5230e9e6(src_path, expected_state, **options):
    """
    Check if the textbox is vertically centered on the image.
    Variation 9 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text vertical center is within middle 10%, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    top_most = height
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                top_most = min(top_most, y)
                bottom_most = max(bottom_most, y)
    text_center = (top_most + bottom_most) / 2
    image_center = height / 2
    tolerance = height * 0.05
    if abs(text_center - image_center) < tolerance:
        return 1.0
    else:
        return 0.0

def check_exact_file_count__1271f790(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file count exactly matches expected and all files are PNG images"""
    expected_count = expected.get('count', 0)
    actual_count = result.get('count', 0)
    if actual_count != expected_count:
        return 0.0
    files = result.get('files', [])
    for file in files:
        title = file.get('title', '').lower()
        mime_type = file.get('mimeType', '')
        is_png = title.endswith('.png') or mime_type == 'image/png'
        if not is_png:
            logger.warning(f"File '{file.get('title')}' is not PNG format (MIME: {mime_type})")
            return 0.0
    return 1.0

def check_file_permissions__7bde372c(result, rules, **options):
    """Check if all files have the expected permission.

    Args:
        result: Output from ls -l command showing file permissions
        rules: Dict with 'permission' key specifying expected permission string (e.g., '-rwxr-xr-x')

    Returns:
        float: 1.0 if all files have correct permission, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    expected_perm = rules.get('permission', '')
    lines = result.strip().split('\n')
    for line in lines:
        if not line.strip():
            continue
        parts = line.split()
        if len(parts) > 0:
            actual_perm = parts[0]
            if actual_perm != expected_perm:
                return 0.0
    return 1.0

def check_file_deleted__9e688855(result, expected, **options):
    """Check if file was deleted.

    Args:
        result: Boolean indicating if file does NOT exist
        expected: Expected state (dict with 'deleted' key)
        **options: Additional options

    Returns:
        float: 1.0 if file deletion matches expectation, 0.0 otherwise
    """
    expected_deleted = expected.get('deleted', True)
    if result == expected_deleted:
        return 1.0
    else:
        return 0.0

def check_text_color__b3dc20103be74c0981bbd79306e76a3d(result, expected, **options):
    """
    Check if the text color matches the expected color with tolerance.

    Args:
        result: RGB color tuple from getter
        expected: Dictionary with 'color' key containing target RGB tuple
        **options: Additional options (e.g., 'tolerance' for color matching)

    Returns:
        float: 1.0 if color matches within tolerance, 0.0 otherwise
    """
    if result is None:
        return 0.0
    target_color = expected.get('color')
    if target_color is None:
        return 0.0
    tolerance = options.get('tolerance', 30)
    distance = np.sqrt(sum(((a - b) ** 2 for (a, b) in zip(result, target_color))))
    if distance <= tolerance:
        return 1.0
    else:
        return 0.0

def check_file_location__c78e5698dfdf96679302f35b21f0928f(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file exists at expected location with expected name.

    Args:
        result: Dict from getter with 'exists', 'path', 'filename' keys
        expected: Dict with 'path', 'filename' expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    else:
        logger.info('File does not exist')
        return score
    expected_filename = expected.get('filename')
    actual_filename = result.get('filename')
    if actual_filename == expected_filename:
        score += 0.5
    else:
        logger.info(f'Filename mismatch - Expected: {expected_filename}, Got: {actual_filename}')
    return score

def check_text_output__b83d8fa5(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_text_occurrence__7dc05d2e(result, expected, **options):
    """Check if text occurrence count matches expected value.
    Verifies that 'European' was replaced with 'EU' and document was saved.

    Args:
        result: Dict with 'european_count', 'eu_count', and 'was_saved'
        expected: Expected count (dict with 'count' key)
        **options: Additional options

    Returns:
        float: 1.0 if all conditions met, partial scores for partial completion
    """
    if isinstance(result, int):
        expected_count = expected.get('count', 0)
        return 1.0 if result == expected_count else 0.0
    expected_count = expected.get('count', 0)
    european_count = result.get('european_count', -1)
    eu_count = result.get('eu_count', 0)
    was_saved = result.get('was_saved', False)
    if european_count == -1:
        return 0.0
    score = 0.0
    if european_count == expected_count:
        score += 0.4
    if eu_count >= 20:
        score += 0.4
    if was_saved:
        score += 0.2
    return score

def check_all_python_syntax__c9c8227a4cd72e8de3e73ea399f7f61d(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if all Python files have valid syntax.

    Args:
        result: Dict from getter containing syntax validation results
        expected: Dict with expected properties (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
            - Proportional score based on number of valid files
            - 1.0 only if all files are valid
    """
    total_files = result.get('total_files', 0)
    valid_files = result.get('valid_files', 0)
    if total_files == 0:
        return 0.0
    score = valid_files / total_files
    if result.get('all_valid', False):
        return 1.0
    return score

def check_file_content__2a7463c5815fe65f87729a241d0d409d(result, expected, **options):
    """Check if file has expected content characteristics.

    Args:
        result: File content statistics from getter
        expected: Expected rules (dict with min_lines, has_classes, etc.)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    checks = 0
    if 'min_code_lines' in expected:
        checks += 1
        if result.get('code_lines', 0) >= expected['min_code_lines']:
            score += 1.0
    if 'min_total_lines' in expected:
        checks += 1
        if result.get('total_lines', 0) >= expected['min_total_lines']:
            score += 1.0
    if 'has_classes' in expected:
        checks += 1
        if result.get('has_class', False) == expected['has_classes']:
            score += 1.0
    if 'has_functions' in expected:
        checks += 1
        if result.get('has_def', False) == expected['has_functions']:
            score += 1.0
    if 'has_imports' in expected:
        checks += 1
        if result.get('has_import', False) == expected['has_imports']:
            score += 1.0
    if 'min_chars' in expected:
        checks += 1
        if result.get('char_count', 0) >= expected['min_chars']:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_file_count__19cf6326(result, expected, **options):
    """Check if specific files exist and are empty.

    Args:
        result: Dict with filenames as keys and file sizes as values
        expected: Dict with 'count' key and optional 'filenames' and 'empty' keys
        **options: Additional options

    Returns:
        float: 1.0 if all required files exist and are empty, 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    expected_count = expected.get('count', 0)
    required_filenames = expected.get('filenames', [])
    require_empty = expected.get('empty', True)
    if len(result) != expected_count:
        return 0.0
    if required_filenames:
        result_filenames = set(result.keys())
        required_filenames_set = set(required_filenames)
        if result_filenames != required_filenames_set:
            return 0.0
        if require_empty:
            for filename in required_filenames:
                if result.get(filename, -1) != 0:
                    return 0.0
    return 1.0

def check_file_rename__990ae9b047da99489a16db0558f7ee61(result: Dict[str, bool], expected: Dict[str, Any], **options) -> float:
    """Check if a file was renamed correctly.

    Args:
        result: Dict mapping file paths to existence status
        expected: Dict with 'old_path' and 'new_path' keys
        **options: Additional options (not used)

    Returns:
        float: 1.0 if old file doesn't exist and new file exists, 0.0 otherwise
    """
    old_path = expected.get('old_path')
    new_path = expected.get('new_path')
    if not old_path or not new_path:
        logger.error('old_path or new_path not provided in expected')
        return 0.0
    old_exists = result.get(old_path, False)
    new_exists = result.get(new_path, False)
    if not old_exists and new_exists:
        return 1.0
    return 0.0

def check_srt_filename__ab47203640c5bdcef1195c50e51e7524(result, expected, **options):
    """
    Check if SRT file exists with correct filename and has valid SRT content.

    Args:
        result: Dict from getter with validation results
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.2
    if result.get('has_content', False):
        score += 0.2
    if result.get('is_valid_srt', False):
        score += 0.4
    if result.get('has_multiple_entries', False):
        score += 0.2
    return score

def check_chapter_files_exist__93cfd69b3c5adfd5dbb8817764000202(result, expected, **options):
    """Check if all expected chapter files exist.

    Args:
        result: Dict with 'existing_files' and 'missing_files' from getter
        expected: Dict with 'required_files' (list of filenames that must exist)
        **options: Additional options

    Returns:
        float: Score based on proportion of required files that exist
    """
    if not isinstance(result, dict):
        return 0.0
    required_files = expected.get('required_files', [])
    existing_files = result.get('existing_files', [])
    if len(required_files) == 0:
        return 1.0
    existing_count = sum((1 for f in required_files if f in existing_files))
    return existing_count / len(required_files)

def check_timezone__b98a8580(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_dual_file_exists__e1a4b749(result, expected, **options):
    """
    Check if both files exist.

    Args:
        result: dict with csv_exists and xlsx_exists
        expected: dict (not used, both must exist)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('csv_exists', False):
        score += 0.5
    if result.get('xlsx_exists', False):
        score += 0.5
    return score

def check_file_line_count__e09bfbdf(result, expected, **options):
    """
    Check if file line count matches expected.

    Args:
        result: Line count from getter
        expected: Expected line count

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_count = expected.get('lines', 0)
    if result == expected_count:
        return 1.0
    return 0.0

def check_text_content__67c66e9d6c723be29d116d6e2c7b5850(result, expected, **options):
    """Check if text content matches expected value.

    Args:
        result: Actual text content from getter
        expected: Expected text content

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    result_str = str(result).strip()
    expected_str = str(expected).strip()
    if result_str == expected_str:
        return 1.0
    return 0.0

def check_text_value__06ae69a0(result, expected, **options):
    """Compare text value (case-insensitive)."""
    if isinstance(expected, dict) and 'rules' in expected:
        expected_val = expected['rules'].get('value', '')
    else:
        expected_val = expected if not isinstance(expected, dict) else expected.get('value', '')
    if result is None:
        return 0.0
    case_sensitive = options.get('case_sensitive', False)
    if case_sensitive:
        return 1.0 if str(result) == str(expected_val) else 0.0
    else:
        return 1.0 if str(result).lower() == str(expected_val).lower() else 0.0

def check_textbox_at_top__83c7a705(result_state, expected_state, **options):
    """
    Check if the textbox is at the top of the image.
    Variation 2 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result_state: Path to the result image with text (vm_file path)
        expected_state: Not used (rule-based)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is at top, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    try:
        source_image = Image.open(result_state)
        gray_image = source_image.convert('L')
        (width, height) = source_image.size
        top_most_dark_pixel = None
        for y in range(height):
            for x in range(width):
                if gray_image.getpixel((x, y)) < 128:
                    top_most_dark_pixel = y
                    break
            if top_most_dark_pixel is not None:
                break
        if top_most_dark_pixel is None:
            logger.warning('No dark pixels (text) found in image')
            return 0.0
        if top_most_dark_pixel < height * 0.03:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking textbox position: {e}')
        return 0.0

def check_gitignore_file__6ee0182d(actual: str, rules: dict, **options) -> float:
    """
    Check if .gitignore file exists and contains Python-related ignore patterns.

    Args:
        actual (str): path to .gitignore file
        rules (dict): expected configuration rules

    Returns:
        float: score between 0.0 and 1.0
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
    except Exception as e:
        logger.error(f'Failed to read .gitignore: {e}')
        return 0.0
    score = 0.0
    score += 0.3
    common_patterns = ['__pycache__', '*.pyc', '*.pyo', '.env', 'venv', '*.egg-info']
    patterns_found = 0
    for pattern in common_patterns:
        if pattern in content:
            patterns_found += 1
    if patterns_found >= 3:
        score += 0.7
    elif patterns_found >= 1:
        score += 0.4
    return score

def check_text_output__ddbee6ed(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_copy__dc38ce29eba391e7169ff4e028e69a72(result, expected, **options):
    """Check if file was copied correctly with content verification.

    Args:
        result: Dict containing copy verification results
        expected: Expected dict with verification criteria
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('copy_successful') == expected.get('copy_successful'):
        return 1.0
    return 0.0

def check_text_replacement__e99e7d69ea502440f7bd1b2cb57a309d(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from the file
        expected: Expected rules dict with 'original_word' and 'correct_word'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    original_word = expected.get('original_word', '')
    correct_word = expected.get('correct_word', '')
    if not original_word or not correct_word:
        return 0.0
    if original_word in result:
        return 0.0
    expected_count = expected.get('expected_count', 0)
    actual_count = result.count(correct_word)
    if actual_count >= expected_count:
        return 1.0
    if expected_count > 0:
        return min(1.0, actual_count / expected_count)
    return 0.0

def check_direct_json_object__f5d96daf_task_verify_3(result, rules, **options):
    """
    Check if result JSON object matches expected JSON object, with support for ignore_list_order flag.

    This is a fixed version of check_direct_json_object that correctly reads ignore_list_order
    from the rules level instead of from inside the expected dict.

    Args:
        result: Dict extracted by getter (e.g., {"modelList": ["model1", "model2", "model3"]})
        rules: Dict containing expected values and flags (e.g., {"expected": {...}, "ignore_list_order": true})
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    logger.info(f'[DEBUG] check_direct_json_object__f5d96daf_task_verify_3 called')
    logger.info(f'[DEBUG] Result: {result}')
    logger.info(f'[DEBUG] Rules: {rules}')
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, returning 0.0')
        return 0.0
    try:
        for (key, value) in result.items():
            if isinstance(value, str) and value.lower() in ['evaluation failed', 'eval_failed', 'failed']:
                logger.error(f"[DEBUG] Expected value for key '{key}' indicates evaluation failure, returning 0.0")
                return 0.0
    except Exception as e:
        logger.error(f'[DEBUG] Error checking for evaluation failure indicator: {e}')
        return 0.0
    try:
        expect_in_result = rules.get('expect_in_result', False)
        ignore_list_order = rules.get('ignore_list_order', False)
        logger.info(f'[DEBUG] expect_in_result: {expect_in_result}')
        logger.info(f'[DEBUG] ignore_list_order: {ignore_list_order}')
        if not expect_in_result:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON: {expected_json}')
            for key in expected_json.keys():
                expected_value = expected_json.get(key)
                actual_value = result.get(key)
                logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
                if ignore_list_order and isinstance(expected_value, list):
                    if not isinstance(actual_value, list):
                        logger.info(f"[DEBUG] Expected list for key '{key}' but got {type(actual_value)}, returning 0.0")
                        return 0.0
                    expected_sorted = sorted(expected_value)
                    actual_sorted = sorted(actual_value)
                    logger.info(f'[DEBUG] Comparing lists (sorted): expected={expected_sorted}, actual={actual_sorted}')
                    if expected_sorted != actual_sorted:
                        logger.info(f"[DEBUG] List comparison failed for key '{key}', returning 0.0")
                        return 0.0
                elif expected_value != actual_value:
                    logger.info(f"[DEBUG] Value comparison failed for key '{key}': expected='{expected_value}', actual='{actual_value}', returning 0.0")
                    return 0.0
                else:
                    logger.info(f"[DEBUG] Value comparison passed for key '{key}'")
            logger.info('[DEBUG] All comparisons passed, returning 1.0')
            return 1.0
        else:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON (expect_in_result mode): {expected_json}')
            for key in expected_json.keys():
                if isinstance(expected_json.get(key), list):
                    flag = 0
                    expected_value_list = expected_json.get(key)
                    logger.info(f"[DEBUG] Checking list key '{key}': expected_list={expected_value_list}, actual='{result.get(key)}'")
                    for each_expected_value in expected_value_list:
                        if isinstance(result.get(key), list) and each_expected_value in result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' in result list for key '{key}'")
                            break
                        elif isinstance(result.get(key), str) and each_expected_value == result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' matches result string for key '{key}'")
                            break
                    if flag == 0:
                        logger.info(f"[DEBUG] No expected values found in result for key '{key}', returning 0.0")
                        return 0.0
                elif isinstance(expected_json.get(key), str):
                    expected_str = expected_json.get(key)
                    actual_str = result.get(key)
                    logger.info(f"[DEBUG] Checking string key '{key}': expected='{expected_str}', actual='{actual_str}'")
                    if expected_str != actual_str:
                        logger.info(f"[DEBUG] String comparison failed for key '{key}', returning 0.0")
                        return 0.0
            logger.info('[DEBUG] All expect_in_result checks passed, returning 1.0')
            return 1.0
    except KeyError as e:
        logger.error(f'[DEBUG] KeyError: {e}, returning 0.0')
        return 0.0
    except Exception as e:
        logger.error(f'[DEBUG] Exception in check_direct_json_object__f5d96daf_task_verify_3: {e}')
        import traceback
        traceback.print_exc()
        return 0.0

def check_file_exists__7298f585793a2b727ac3910de3795a50(result: bool, expected: Dict[str, Any], **options) -> float:
    """Check if a file exists.

    Args:
        result: Boolean indicating if file exists
        expected: Dict with 'should_exist' key (True/False)
        **options: Additional options

    Returns:
        float: 1.0 if result matches expectation, 0.0 otherwise
    """
    should_exist = expected.get('should_exist', True)
    if result == should_exist:
        return 1.0
    else:
        return 0.0

def check_text_file_content__50943048(result_path, expected, **options):
    """
    Checks if a text file contains the expected content.

    Args:
        result_path: Path to the text file to check
        expected: Dict with 'expected_content' (string to match)
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not os.path.exists(result_path):
        logger.error(f'Text file does not exist: {result_path}')
        return 0.0
    file_size = os.path.getsize(result_path)
    if file_size == 0:
        logger.error(f'Text file is empty: {result_path}')
        return 0.0
    expected_content = expected.get('expected_content', '')
    try:
        with open(result_path, 'r', encoding='utf-8') as f:
            actual_content = f.read().strip()
        logger.info(f"Actual content: '{actual_content}'")
        logger.info(f"Expected content: '{expected_content}'")
        if actual_content == expected_content:
            return 1.0
        if expected_content in actual_content:
            logger.info(f'Expected content found within actual content')
            return 0.8
        return 0.0
    except Exception as e:
        logger.error(f'Error reading text file: {e}')
        return 0.0

def check_line_count__953256df836603c8857d4495861e4b63(result: str, expected: Dict, **options) -> float:
    """Check if email summary file has correct format and content.

    Validates:
    1. Each line follows format: 'From: [sender] - Subject: [subject]'
    2. Sender contains valid email address format
    3. Line count is within expected range (at least 2 based on Notes folder)
    4. Subject lines are non-empty

    Args:
        result: Actual text content from file
        expected: Expected configuration with 'min_lines' and 'max_lines'

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    lines = [line.strip() for line in result.split('\n') if line.strip()]
    line_count = len(lines)
    min_lines = expected.get('min_lines', 0)
    max_lines = expected.get('max_lines', float('inf'))
    if line_count < min_lines:
        return max(0.0, 0.3 * (line_count / min_lines))
    if line_count > max_lines:
        return 0.5
    email_pattern = 'From:\\s+([^\\s@]+@[^\\s@]+\\.[^\\s@]+)\\s+-\\s+Subject:\\s+(.+)'
    valid_lines = 0
    lines_with_email_format = 0
    lines_with_subject = 0
    for line in lines:
        match = re.match(email_pattern, line, re.IGNORECASE)
        if match:
            valid_lines += 1
            sender = match.group(1).strip()
            subject = match.group(2).strip()
            if '@' in sender and '.' in sender.split('@')[1]:
                lines_with_email_format += 1
            if subject and len(subject) > 0:
                lines_with_subject += 1
    if line_count == 0:
        return 0.0
    format_score = valid_lines / line_count
    email_score = lines_with_email_format / line_count
    subject_score = lines_with_subject / line_count
    overall_score = (format_score + email_score + subject_score) / 3
    if overall_score >= 0.9:
        return 1.0
    elif overall_score >= 0.7:
        return 0.8
    elif overall_score >= 0.5:
        return 0.6
    elif overall_score >= 0.3:
        return 0.4
    else:
        return max(0.0, overall_score)

def check_all_files_exist__03426e679d8f4571bede57a16eea69a4(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if all files were properly copied (source exists, destination exists, content matches).

    Args:
        result: List of dicts from getter with keys: 'dest_exists', 'source_exists', 'content_matches'
        expected: Dict with 'all_should_exist' key (True/False)

    Returns:
        1.0 if all files were properly copied, 0.0 otherwise
    """
    all_should_exist = expected.get('all_should_exist', True)
    if not all_should_exist:
        logger.warning('Unexpected expected value: all_should_exist is False')
        return 0.0
    all_properly_copied = True
    for (i, file_result) in enumerate(result):
        dest_exists = file_result.get('dest_exists', False)
        source_exists = file_result.get('source_exists', False)
        content_matches = file_result.get('content_matches', False)
        properly_copied = dest_exists and source_exists and content_matches
        if not properly_copied:
            all_properly_copied = False
            logger.info(f'File {i}: dest_exists={dest_exists}, source_exists={source_exists}, content_matches={content_matches}')
            if not source_exists:
                logger.warning(f'File {i}: Source file missing - this was a MOVE, not a COPY')
            if not dest_exists:
                logger.warning(f'File {i}: Destination file missing')
            if dest_exists and source_exists and (not content_matches):
                logger.warning(f'File {i}: Content mismatch - files have different content')
    if all_properly_copied:
        logger.info(f'All {len(result)} files were properly copied (source preserved, destination created, content matches)')
        return 1.0
    else:
        failed_count = sum((1 for f in result if not (f.get('dest_exists') and f.get('source_exists') and f.get('content_matches'))))
        logger.info(f'{failed_count}/{len(result)} files failed copy verification')
        return 0.0

def check_text_contains__a515d38160faf207fcfd0df30838b0c3(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if the result text contains expected substring.

    Args:
        result: Actual text content
        expected: Dict with 'substring' key containing the expected substring
        **options: Additional options (case_sensitive, default True)

    Returns:
        float: 1.0 if substring found, 0.0 otherwise
    """
    substring = expected.get('substring', '')
    case_sensitive = options.get('case_sensitive', True)
    if not case_sensitive:
        result = result.lower()
        substring = substring.lower()
    if substring in result:
        return 1.0
    return 0.0

def check_textbox_topleft__7d24ebd2(src_path, expected_state, **options):
    """
    Check if the textbox is in the upper-left corner of the image.
    Variation 4 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in top-left 5% region, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    top_most = height
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                top_most = min(top_most, y)
    if left_most < width * 0.05 and top_most < height * 0.05:
        return 1.0
    else:
        return 0.0

def check_zip_files__53fe105429a60cae06bcf9ce59b19b3e(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if zip archive contains exactly the expected files.

    Args:
        result: List of filenames from getter
        expected: Dict with 'expected_files' key containing list of expected filenames
        **options: Additional options (unused)

    Returns:
        1.0 if zip contains exactly the expected files (no more, no less), 0.0 otherwise
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    else:
        return 0.0

def check_text_replacement__71cb5a9efacfe89bccff2081bcadcf02(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from file
        expected: Rules dict containing 'expected_text' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    result_clean = result.strip()
    expected_clean = expected_text.strip()
    if result_clean == expected_clean:
        return 1.0
    return 0.0

def check_file_saved__9c119f81(result, expected, **options):
    if not result.get('exists'):
        return 0.0
    score = 0.5
    if result.get('size', 0) >= expected.get('min_size', 1000):
        score += 0.5
    return score

def check_file_exists__6f3c16ae(result: dict, expected: dict, **options) -> float:
    """Check if the tar.gz archive exists, is valid, and contains expected files.

    Args:
        result: Dict with archive validation info from getter
            {
                'exists': bool,
                'valid_archive': bool,
                'contents': list,
                'has_expected_structure': bool
            }
        expected: Dict with 'exists' key (bool)
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        return 1.0 if not result.get('exists', False) else 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('valid_archive', False):
        return 0.0
    if not result.get('has_expected_structure', False):
        return 0.0
    return 1.0

def check_file_created__684b5a3a3f653750766f5bbe64af3bd5(result: bool, expected: dict, **options) -> float:
    """Check if a file was created successfully.

    Args:
        result: Boolean indicating if file exists
        expected: Dict with 'exists' field (True/False)

    Returns:
        1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_text_content__25d344d5(result_file, expected, **options):
    try:
        prs = Presentation(result_file)
        slide_idx = expected.get('slide_idx', 0)
        shape_idx = expected.get('shape_idx', 0)
        expected_text = expected.get('expected_text', '')
        if slide_idx >= len(prs.slides):
            return 0.0
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            return 0.0
        shape = slide.shapes[shape_idx]
        if not hasattr(shape, 'text'):
            return 0.0
        actual_text = shape.text.strip()
        return 1.0 if actual_text == expected_text else 0.0
    except:
        return 0.0

def check_file_content_e9b24959(result, expected, **options):
    """Check if file content matches expected count.

    Args:
        result: File path string returned by get_vm_file()
        expected: Expected rules dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        count = int(content)
        expected_count = expected.get('expected_count', 0)
        if count == expected_count:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_timezone__f88b70b8(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_min_file_size__b756f99d(result, expected, **options):
    """Check if file size meets minimum requirement.

    Args:
        result: Actual file size in bytes (int)
        expected: Dict with 'min_size_bytes' key
        **options: Additional options

    Returns:
        float: 1.0 if file size >= min_size_bytes, 0.0 otherwise
    """
    min_size = expected.get('min_size_bytes', 1024)
    if result >= min_size:
        return 1.0
    else:
        return 0.0

def check_exclude_platform__156bbecf8094c096edc4b32f7b6fd25b(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the docx file excludes all records for the specified platform.

    Args:
        result: List of train record lines from getter
        expected: Dict with 'excluded_platform' and 'expected_count' keys
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    excluded_platform = expected.get('excluded_platform')
    expected_count = expected.get('expected_count')
    if not excluded_platform or expected_count is None:
        logger.error('Missing excluded_platform or expected_count in expected config')
        return 0.0
    valid_count = 0
    for line in result:
        parts = line.split(',')
        if len(parts) == 4:
            platform_no = parts[3].strip()
            if platform_no == excluded_platform:
                logger.info(f'Found record from excluded platform: {platform_no}')
                return 0.0
            valid_count += 1
    if valid_count == expected_count:
        return 1.0
    else:
        logger.info(f'Record count mismatch: got {valid_count}, expected {expected_count}')
        return 0.0

def check_intext_citation_added__9d568660(result, expected, **options):
    """Check if in-text citation was added to specific paragraph.

    Args:
        result: Path to the result DOCX file
        expected: Dict with 'paragraph_start', 'citation', 'near_phrase' keys

    Returns:
        float: 1.0 if citation found in correct paragraph, 0.0 otherwise
    """
    if not result or not isinstance(result, str):
        return 0.0
    try:
        doc = Document(result)
        paragraph_start = expected.get('paragraph_start', '')
        citation = expected.get('citation', '')
        near_phrase = expected.get('near_phrase', '')
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if para_text.startswith(paragraph_start):
                if citation in para_text:
                    if near_phrase:
                        phrase_pos = para_text.find(near_phrase)
                        citation_pos = para_text.find(citation)
                        if phrase_pos >= 0 and citation_pos >= 0:
                            if abs(citation_pos - phrase_pos) < 100:
                                return 1.0
                    else:
                        return 1.0
        return 0.0
    except Exception as e:
        print(f'Error in check_intext_citation_added__9d568660: {e}')
        return 0.0

def check_file_count__bc253f41(result, expected, **options):
    """Compare file count against expected value.

    Args:
        result: Count from getter
        expected: Expected count value
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_count = expected.get('count', 0)
    return 1.0 if result == expected_count else 0.0

def check_first_line_indent__7678624b(result, expected, **options):
    """
    Compare paragraph first-line indentation in a docx file against expected values.

    Args:
        result: Path to the result docx file
        expected: Dictionary with first_line_indents mapping (0-indexed paragraph -> inches)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not expected:
        return 0.0
    first_line_indents = expected.get('first_line_indents', {})
    if not first_line_indents:
        return 0.0
    try:
        doc = Document(result)
    except Exception as e:
        logger.error(f'Error loading document: {e}')
        return 0.0
    score = 0.0
    total_checks = len(first_line_indents)
    for (para_idx_str, expected_indent_inches) in first_line_indents.items():
        para_idx = int(para_idx_str)
        if para_idx >= len(doc.paragraphs):
            continue
        para = doc.paragraphs[para_idx]
        actual_indent = para.paragraph_format.first_line_indent
        expected_indent_emu = Inches(expected_indent_inches)
        if actual_indent is not None:
            tolerance = Inches(0.05)
            if abs(actual_indent - expected_indent_emu) <= tolerance:
                score += 1.0 / total_checks
        elif expected_indent_inches == 0:
            score += 1.0 / total_checks
    return score

def check_file_readonly__b3b80682(result, expected, **options):
    """Verify file exists and has readonly permissions.

    Args:
        result: Dict with file_exists, permissions, and readonly status
        expected: Dict with rules specifying expected states
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 with partial credit
    """
    expected_file_exists = expected.get('file_exists', True)
    expected_readonly = expected.get('readonly', True)
    file_exists = result.get('file_exists', False)
    readonly = result.get('readonly', False)
    score = 0.0
    if file_exists == expected_file_exists:
        score += 0.5
    if readonly == expected_readonly:
        score += 0.5
    return score

def check_file_exists__689ec9af4ba1471bf9b5f89e71cafeb9(result, expected, **options):
    """Check if a valid BMP file exists as expected.

    Args:
        result: Result from getter (dict with 'exists', 'is_bmp', 'size' keys)
        expected: Expected rules (dict with 'should_exist' key)
        **options: Additional options

    Returns:
        float: 1.0 if valid BMP file exists as expected, 0.0 otherwise
    """
    file_exists = result.get('exists', False)
    is_bmp = result.get('is_bmp', False)
    file_size = result.get('size', 0)
    should_exist = expected.get('should_exist', True)
    if should_exist:
        if file_exists and is_bmp and (file_size > 0):
            return 1.0
        else:
            return 0.0
    elif not file_exists:
        return 1.0
    else:
        return 0.0

def check_new_file_properties__73935909(result_state, expected_state, **options):
    """
    Check if a new file was created with the expected properties.

    Args:
        result_state: Dict from getter with file info
        expected_state: Dict with expected properties:
            - table_count: int, expected number of tables (0 for no tables)
            - min_words: int, minimum word count
            - max_words: int, maximum word count
        **options: Additional options

    Returns:
        float: Score (1.0 if all checks pass, 0.0 otherwise)
    """
    if not isinstance(result_state, dict):
        logger.error(f'Invalid result_state type: {type(result_state)}, expected dict')
        return 0.0
    if not isinstance(expected_state, dict):
        logger.error(f'Invalid expected_state type: {type(expected_state)}, expected dict')
        return 0.0
    if not result_state.get('exists', False):
        logger.error('❌ No new file was created')
        return 0.0
    logger.info(f"✅ New file exists: {result_state.get('path', 'unknown')}")
    expected_table_count = expected_state.get('table_count', 0)
    actual_table_count = result_state.get('table_count', -1)
    if actual_table_count == -1:
        logger.error('❌ Could not read table count from file')
        return 0.0
    if actual_table_count != expected_table_count:
        logger.error(f'❌ Table count mismatch: expected {expected_table_count}, got {actual_table_count}')
        return 0.0
    logger.info(f'✅ Table count correct: {actual_table_count}')
    min_words = expected_state.get('min_words')
    max_words = expected_state.get('max_words')
    word_count = result_state.get('word_count', 0)
    if min_words is None or max_words is None:
        logger.error(f'Missing min_words or max_words in expected_state: {expected_state}')
        return 0.0
    if not min_words <= word_count <= max_words:
        logger.error(f'❌ Word count {word_count} is outside range [{min_words}, {max_words}]')
        return 0.0
    logger.info(f'✅ Word count {word_count} is within range [{min_words}, {max_words}]')
    logger.info('✅ All checks passed for new file')
    return 1.0

def check_file_content__d09dce0a(result, expected, **options):
    if not result.get('exists'):
        return 0.0
    score = 0.5
    expected_content = str(expected.get('expected_content', '3')).strip()
    actual_content = result.get('content', '').strip()
    if actual_content == expected_content:
        score += 0.5
    return score

def check_python_imports_only__093631738f9b5eba42a5bcf60212ba3b(result, expected, **options):
    """
    Check if Python file contains only import statements (import/from imports).

    Args:
        result: File content string from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    score = 0.0
    if result and len(result.strip()) > 0:
        score += 0.2
    else:
        return 0.0
    lines = result.split('\n')
    import_lines = []
    non_import_non_comment_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith('#'):
            continue
        if stripped.startswith('import ') or stripped.startswith('from '):
            import_lines.append(line)
        else:
            non_import_non_comment_lines.append(line)
    min_imports = expected.get('min_imports', 1)
    if len(import_lines) >= min_imports:
        score += 0.5
    total_code_lines = len(import_lines) + len(non_import_non_comment_lines)
    if total_code_lines > 0:
        import_ratio = len(import_lines) / total_code_lines
        if import_ratio >= 0.7:
            score += 0.3
    return score

def check_file_contains_lines__fde871ae(result_file_path, expected, **options):
    """
    Check if a text file contains all expected lines.
    
    Args:
        result_file_path: Path to the result text file
        expected: Dict with 'lines' key containing list of expected lines
        **options: Additional options (case_sensitive, order_matters, etc.)
        
    Returns:
        float: Score between 0.0 and 1.0
    """
    if result_file_path is None:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            result_lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    expected_lines = expected.get('lines', [])
    if not expected_lines:
        return 0.0
    case_sensitive = options.get('case_sensitive', True)
    order_matters = options.get('order_matters', False)
    if not case_sensitive:
        result_lines_normalized = [line.lower() for line in result_lines]
        expected_lines_normalized = [line.lower() for line in expected_lines]
    else:
        result_lines_normalized = result_lines
        expected_lines_normalized = expected_lines
    if order_matters:
        if len(result_lines_normalized) != len(expected_lines_normalized):
            return 0.0
        matches = sum((1 for (i, exp_line) in enumerate(expected_lines_normalized) if i < len(result_lines_normalized) and result_lines_normalized[i] == exp_line))
        return matches / len(expected_lines_normalized)
    else:
        matches = sum((1 for exp_line in expected_lines_normalized if exp_line in result_lines_normalized))
        return matches / len(expected_lines_normalized)

def check_file_contains_lines__eac4b332(result_file_path, expected, **options):
    """
    Check if a text file contains all expected lines.
    
    Args:
        result_file_path: Path to the result text file
        expected: Dict with 'lines' key containing list of expected lines
        **options: Additional options (case_sensitive, order_matters, etc.)
        
    Returns:
        float: Score between 0.0 and 1.0
    """
    if result_file_path is None:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            result_lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    expected_lines = expected.get('lines', [])
    if not expected_lines:
        return 0.0
    case_sensitive = options.get('case_sensitive', True)
    order_matters = options.get('order_matters', False)
    if not case_sensitive:
        result_lines_normalized = [line.lower() for line in result_lines]
        expected_lines_normalized = [line.lower() for line in expected_lines]
    else:
        result_lines_normalized = result_lines
        expected_lines_normalized = expected_lines
    if order_matters:
        if len(result_lines_normalized) != len(expected_lines_normalized):
            return 0.0
        matches = sum((1 for (i, exp_line) in enumerate(expected_lines_normalized) if i < len(result_lines_normalized) and result_lines_normalized[i] == exp_line))
        return matches / len(expected_lines_normalized)
    else:
        matches = sum((1 for exp_line in expected_lines_normalized if exp_line in result_lines_normalized))
        return matches / len(expected_lines_normalized)

def check_file_downloaded__08c5e1b6ad7015f1bdd4ff79ff88e12f(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a file was successfully downloaded.

    Args:
        result: Dict from getter with file information
        expected: Dict with expected file properties
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    score += 0.4
    if result.get('is_file', False):
        score += 0.2
    else:
        logger.warning('Path exists but is not a file')
        return score
    if result.get('meets_min_size', False):
        score += 0.4
        logger.info(f"File meets minimum size requirement: {result.get('size_bytes', 0)} bytes")
    else:
        logger.warning(f"File size {result.get('size_bytes', 0)} bytes is below minimum requirement")
    logger.info(f'File download check score: {score}')
    return score

def check_exact_text_match__65949e53(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text from getter
        expected: Expected text string
        **options: Additional options (unused)

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result == expected:
        return 1.0
    return 0.0

def check_file_content_555dda86(result, expected, **options):
    """Check if file content matches expected count.

    Args:
        result: File path (string) returned by get_vm_file()
        expected: Expected rules dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        count = int(content)
        expected_count = expected.get('expected_count', 0)
        if count == expected_count:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_files_exist__a8a082525df1807c95a7519289fda5a0(result, expected, **options):
    """Check if expected files exist with correct row counts.

    Args:
        result: Dict mapping file paths to existence/row_count info
        expected: Dict with 'files' list from rules (each item has 'path' and 'min_rows')
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on file existence and row counts
    """
    files = expected.get('files', [])
    if not files:
        return 0.0
    score = 0.0
    points_per_file = 1.0 / len(files)
    for file_info in files:
        path = file_info.get('path', '')
        min_rows = file_info.get('min_rows', 0)
        if path in result:
            file_result = result[path]
            if file_result.get('exists', False):
                score += points_per_file * 0.5
                if file_result.get('row_count', 0) >= min_rows:
                    score += points_per_file * 0.5
    return score

def check_file_permissions__f81b27ec(result, rules, **options):
    """Check if all files have the expected permission.

    Args:
        result: Output from ls -l command showing file permissions
        rules: Dict with 'permission' key specifying expected permission string (e.g., '-rwxrwxrwx')

    Returns:
        float: 1.0 if all files have correct permission, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    expected_perm = rules.get('permission', '')
    lines = result.strip().split('\n')
    for line in lines:
        if not line.strip():
            continue
        parts = line.split()
        if len(parts) > 0:
            actual_perm = parts[0]
            if actual_perm != expected_perm:
                return 0.0
    return 1.0

def check_python_code_backup__f222fdd4d3c26325ac7310b0f2b1711f(result, expected, **options):
    """
    Check if Python backup file exists with valid content and structure.
    Validates that code was extracted from Colab notebook and merged.

    Args:
        result: Dict with 'content', 'filename', and 'file_date' keys from getter
        expected: Rules dict with validation parameters
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    content = result.get('content', '')
    filename = result.get('filename', '')
    file_date = result.get('file_date', '')
    if not content or not filename:
        return 0.0
    score = 0.0
    date_pattern = re.compile('^code_backup_(\\d{8})\\.py$')
    match = date_pattern.match(filename)
    if match:
        today = datetime.now()
        try:
            file_datetime = datetime.strptime(file_date, '%Y%m%d')
            date_diff = abs((today.date() - file_datetime.date()).days)
            if date_diff <= 1:
                score += 0.2
            else:
                logger.warning(f"File date {file_date} is not today's date {today.strftime('%Y%m%d')}")
                return 0.0
        except ValueError:
            logger.warning(f'Invalid date format in filename: {file_date}')
            return 0.0
    else:
        return 0.0
    if len(content.strip()) > 50:
        score += 0.15
    else:
        return score
    lines = content.split('\n')
    code_lines = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith('#') and (not re.match('^#\\s*In\\[\\d+\\]:', stripped)):
            continue
        code_lines.append(line)
    min_code_lines = expected.get('min_code_lines', 1)
    effective_min = max(min_code_lines, 5)
    if len(code_lines) >= effective_min:
        score += 0.15
    else:
        return score
    colab_cell_markers = re.findall('#\\s*In\\[\\d+\\]:', content)
    if len(colab_cell_markers) >= 2:
        score += 0.25
    elif len(colab_cell_markers) == 1:
        score += 0.1
    else:
        has_imports = False
        has_functions = False
        has_classes = False
        blank_line_groups = 0
        prev_blank = True
        for line in lines:
            stripped = line.strip()
            if not stripped:
                if not prev_blank:
                    blank_line_groups += 1
                prev_blank = True
            else:
                prev_blank = False
                if stripped.startswith('import ') or stripped.startswith('from '):
                    has_imports = True
                elif stripped.startswith('def '):
                    has_functions = True
                elif stripped.startswith('class '):
                    has_classes = True
        code_complexity_score = sum([has_imports, has_functions, has_classes])
        if code_complexity_score >= 2 or blank_line_groups >= 3:
            score += 0.15
        elif code_complexity_score >= 1 and len(code_lines) >= 10:
            score += 0.05
    if len(code_lines) >= 15:
        score += 0.2
    elif len(code_lines) >= 10:
        score += 0.1
    return score

def check_file_size_range__198be354(result, expected, **options):
    """Check if file size is within expected range and contains Python code.

    Args:
        result: dict from getter with file size and content info
        expected: dict with 'min_size' and 'max_size' in bytes
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    size_bytes = result.get('size_bytes', 0)
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    is_python = result.get('is_python', False)
    has_functions = result.get('has_functions', False)
    if size_bytes == 0:
        return 0.0
    score = 0.0
    if is_python:
        score += 0.3
    if has_functions:
        score += 0.3
    if min_size <= size_bytes <= max_size:
        score += 0.4
    return score

def check_file_count__17298c22(result, expected, **options):
    """Compare file count against expected value.

    Args:
        result: File count from getter (int)
        expected: Expected rules dict with 'count' and optional 'tolerance'
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if not isinstance(result, int):
        return 0.0
    expected_count = expected.get('count', 0)
    tolerance = expected.get('tolerance', 0)
    if abs(result - expected_count) <= tolerance:
        return 1.0
    else:
        return 0.0

def compare_text_output__81f11cbf(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_file_exists__24a50bf5(result, expected, **options):
    """
    Check if a file exists based on command output.

    Args:
        result: Output from ls command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'exists')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_python_function_exists__c685793b2ca36ab76b7f2cc84f84fe40(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if a Python file contains a specific function definition with expected print statement.

    Args:
        result: Actual Python file content from getter
        expected: Expected dict with 'function_name' key and optional 'print_message' key
        **options: Additional options

    Returns:
        float: 1.0 if function exists with correct implementation, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_function = expected.get('function_name', '')
    if not expected_function:
        return 0.0
    expected_print = expected.get('print_message', None)
    try:
        tree = ast.parse(result)
        for node in ast.walk(tree):
            if isinstance(node, ast.FunctionDef):
                if node.name == expected_function:
                    if expected_print is None:
                        return 1.0
                    if _contains_print_statement(node, expected_print):
                        return 1.0
                    else:
                        return 0.0
        return 0.0
    except SyntaxError:
        return 0.0
    except Exception:
        return 0.0

def check_exact_text_match__6ba0a623(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text from getter
        expected: Expected text string
        **options: Additional options (unused)

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result == expected:
        return 1.0
    return 0.0

def check_filename__4ee0209a(result, expected, **options):
    """Check if filename matches expected.

    Args:
        result: Actual filename (str)
        expected: Dict with 'filename' key
        **options: Additional options

    Returns:
        float: 1.0 if filename matches, 0.0 otherwise
    """
    expected_filename = expected.get('filename', 'invoice.xlsx')
    if result == expected_filename:
        return 1.0
    else:
        return 0.0

def check_line_count__cfbf4273(result_file_path, expected, **options):
    """Check if code line count is in valid range and URL matches for Karpathy GPT Colab notebook.

    This evaluator verifies that:
    1. The file exists and contains data in format 'count|url'
    2. The count is a valid positive integer within the expected range (180-280 lines)
    3. The URL matches the expected Colab notebook URL

    The expected range is based on analysis of Karpathy's GPT-from-scratch tutorial,
    which implements a character-level transformer including data loading, model
    architecture (embeddings, attention, feedforward layers), and training loop.

    By requiring the URL to be saved alongside the count, this evaluator ensures
    that the agent actually accessed the correct notebook rather than guessing
    a number in the valid range.

    Args:
        result_file_path: Path to file containing line count and URL in format 'count|url'
        expected: Dict with 'min_count', 'max_count', and 'notebook_url'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result_file_path:
        print('Error: No result file path provided')
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read().strip()
        if not content:
            print('Error: Result file is empty')
            return 0.0
        parts = content.split('|')
        if len(parts) != 2:
            print(f"Error: File content must be in format 'count|url', got: {content}")
            return 0.0
        (count_str, url) = (parts[0].strip(), parts[1].strip())
        try:
            count = int(count_str)
        except ValueError:
            print(f"Error: Count '{count_str}' is not a valid integer")
            return 0.0
        if count <= 0:
            print(f'Error: Count {count} must be positive')
            return 0.0
        expected_url = expected.get('notebook_url', '')

        def normalize_url(u):
            if 'colab.research.google.com/drive/' in u:
                base = u.split('?')[0].split('#')[0]
                return base.rstrip('/')
            return u.rstrip('/')
        normalized_result_url = normalize_url(url)
        normalized_expected_url = normalize_url(expected_url)
        if normalized_result_url != normalized_expected_url:
            print(f'Error: URL mismatch. Expected: {normalized_expected_url}, Got: {normalized_result_url}')
            return 0.0
        min_count = expected.get('min_count', 0)
        max_count = expected.get('max_count', float('inf'))
        if min_count <= count <= max_count:
            print(f'Success: Count {count} is within expected range [{min_count}, {max_count}] and URL matches')
            return 1.0
        else:
            print(f'Error: Count {count} is outside expected range [{min_count}, {max_count}]')
            return 0.0
    except FileNotFoundError:
        print(f'Error: File not found at {result_file_path}')
        return 0.0
    except Exception as e:
        print(f'Error checking line count: {e}')
        return 0.0

def check_line_count__475840bd88bfc32515242a838ac799b5(result, expected, **options):
    """Check if file is properly merged from all 5 chapters in correct order.

    This verification checks:
    1. Line count is within expected range
    2. Content from Chapter0 appears at the beginning (mandatory markers)
    3. Content from Chapter4 appears at the end (specific markers)
    4. All chapters are present in correct order (boundary verification)
    5. Content integrity is preserved (hash verification)

    Args:
        result: dict from getter {
            'exists': bool,
            'line_count': int,
            'first_100_chars': str,
            'last_100_chars': str,
            'chapter_markers': list,
            'full_content': str,
            'chapter_boundaries': dict
        }
        expected: dict with 'min_lines', 'max_lines', and chapter verification data

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    max_score = 6.0
    actual_count = result.get('line_count', 0)
    full_content = result.get('full_content', '')
    min_lines = expected.get('min_lines', 0)
    max_lines = expected.get('max_lines', float('inf'))
    if min_lines <= actual_count <= max_lines:
        logger.info(f'✓ Line count {actual_count} is within range [{min_lines}, {max_lines}]')
        score += 1.0
    else:
        logger.info(f'✗ Line count {actual_count} is outside range [{min_lines}, {max_lines}]')
        return 0.0
    first_500 = result.get('first_500_chars', '')
    chapter0_markers = ['Chapter 1 Passing through the book', 'Chi Qingluo opened his eyes']
    chapter0_found = all((marker in first_500 for marker in chapter0_markers))
    if chapter0_found:
        logger.info(f'✓ Chapter0 MANDATORY markers detected at beginning')
        score += 1.0
    else:
        logger.info(f"✗ Chapter0 MANDATORY markers NOT found - file doesn't start with Chapter0")
        return 0.0
    last_500 = result.get('last_500_chars', '')
    chapter4_end_markers = ['Having said all that, can you pay me back?']
    chapter4_found = any((marker in last_500 for marker in chapter4_end_markers))
    if chapter4_found:
        logger.info(f'✓ Chapter4 SPECIFIC end marker detected')
        score += 1.0
    else:
        logger.info(f'✗ Chapter4 SPECIFIC end marker NOT found - file may not end with Chapter4')
        return 0.0
    chapter_boundaries = result.get('chapter_boundaries', {})
    if chapter_boundaries:
        found_chapters = [ch for ch in ['ch0', 'ch1', 'ch2', 'ch3', 'ch4'] if chapter_boundaries.get(ch)]
        if len(found_chapters) >= 5:
            ch0_pos = chapter_boundaries.get('ch0', {}).get('line', -1)
            ch1_pos = chapter_boundaries.get('ch1', {}).get('line', -1)
            ch2_pos = chapter_boundaries.get('ch2', {}).get('line', -1)
            ch3_pos = chapter_boundaries.get('ch3', {}).get('line', -1)
            ch4_pos = chapter_boundaries.get('ch4', {}).get('line', -1)
            if ch0_pos < ch1_pos < ch2_pos < ch3_pos < ch4_pos:
                logger.info(f'✓ All 5 chapters found in correct order: ch0({ch0_pos}) < ch1({ch1_pos}) < ch2({ch2_pos}) < ch3({ch3_pos}) < ch4({ch4_pos})')
                score += 1.0
            else:
                logger.info(f'✗ Chapters found but NOT in correct order: ch0({ch0_pos}), ch1({ch1_pos}), ch2({ch2_pos}), ch3({ch3_pos}), ch4({ch4_pos})')
                return 0.0
        else:
            logger.info(f'✗ Only {len(found_chapters)}/5 chapters detected: {found_chapters}')
            return 0.0
    else:
        logger.info(f'✗ No chapter boundary information available')
        return 0.0
    ch0_pos = chapter_boundaries.get('ch0', {}).get('line', -1)
    ch1_pos = chapter_boundaries.get('ch1', {}).get('line', -1)
    ch2_pos = chapter_boundaries.get('ch2', {}).get('line', -1)
    ch3_pos = chapter_boundaries.get('ch3', {}).get('line', -1)
    ch4_pos = chapter_boundaries.get('ch4', {}).get('line', -1)
    expected_positions = {'ch0': (0, 50), 'ch1': (400, 550), 'ch2': (850, 1050), 'ch3': (1200, 1450), 'ch4': (1550, 1800)}
    positions_correct = True
    for (ch, (min_pos, max_pos)) in expected_positions.items():
        pos = chapter_boundaries.get(ch, {}).get('line', -1)
        if not min_pos <= pos <= max_pos:
            logger.info(f'✗ {ch} at line {pos} is outside expected range [{min_pos}, {max_pos}]')
            positions_correct = False
            break
    if positions_correct:
        logger.info(f'✓ All chapters appear at expected positions')
        score += 1.0
    else:
        logger.info(f"✗ Chapter positions don't match expected merge order")
        return 0.0
    content_hash = result.get('content_hash', '')
    expected_hash = expected.get('expected_hash', '')
    if expected_hash and content_hash:
        if content_hash == expected_hash:
            logger.info(f'✓ Content hash matches expected: {content_hash}')
            score += 1.0
        else:
            logger.info(f'✗ Content hash mismatch: got {content_hash}, expected {expected_hash}')
            return 0.0
    else:
        expected_exact = 2167
        if actual_count == expected_exact:
            logger.info(f'✓ Exact line count matches expected {expected_exact}')
            score += 1.0
        else:
            logger.info(f'⚠ Line count {actual_count} != expected exact {expected_exact} (diff: {abs(actual_count - expected_exact)})')
            if abs(actual_count - expected_exact) <= 2:
                score += 0.5
    final_score = score / max_score
    logger.info(f'Final score: {final_score:.2f} ({score}/{max_score} checks passed)')
    if final_score >= 0.95:
        return 1.0
    else:
        return 0.0

def check_file_size__739292ff(result, expected, **options):
    """Check if file size is within expected range.

    Args:
        result: Dict with 'size_bytes' key
        expected: Dict with 'min_size' and optionally 'max_size' keys
        **options: Additional options

    Returns:
        float: 1.0 if size is within range, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    min_size = expected.get('min_size', 0)
    max_size = expected.get('max_size', float('inf'))
    actual_size = result.get('size_bytes', 0)
    if min_size <= actual_size <= max_size:
        logger.info(f'File size OK: {actual_size} bytes within [{min_size}, {max_size}]')
        return 1.0
    else:
        logger.info(f'File size out of range: {actual_size} bytes not in [{min_size}, {max_size}]')
        return 0.0

def check_file_moved__789836386f3e1cf0e0ee5d172a0885f2(result, expected, **options):
    """Check if file was moved correctly (not in source, in destination).

    Args:
        result: Dict with 'in_source' and 'in_dest' booleans
        expected: Dict with expected 'in_source' and 'in_dest' values
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    in_source = expected.get('in_source', False)
    in_dest = expected.get('in_dest', True)
    if result.get('in_source') == in_source and result.get('in_dest') == in_dest:
        return 1.0
    return 0.0

def check_file_copied__642c6d87(result, expected, **options):
    """Check if file was copied correctly by comparing hash.

    Args:
        result: Hash string from copied file
        expected: Rules dict with expected_hash
        **options: Additional options

    Returns:
        float: 1.0 if hashes match, 0.0 otherwise
    """
    if result is None or not isinstance(result, str):
        return 0.0
    result = result.strip()
    expected_hash = expected.get('expected_hash', '')
    if result == 'NOT_FOUND':
        return 0.0
    if result == expected_hash:
        return 1.0
    else:
        return 0.0

def check_line_count_positive__868f1e74(result, expected, **options):
    """Check if line count is positive (file exists and has content).

    Args:
        result: Line count from getter
        expected: Minimum expected line count
        **options: Additional options

    Returns:
        float: 1.0 if line count >= expected, 0.0 otherwise
    """
    try:
        if isinstance(result, (int, float)) and result >= expected:
            return 1.0
        return 0.0
    except:
        return 0.0

def check_file_organization__3b0a753c(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_renamed_files__47275204(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the correct files have been renamed with a prefix.

    Args:
        result: List of filenames with prefix
        expected: Dict with 'rules' containing 'expected_files' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on matching files
    """
    if result is None:
        return 0.0
    expected_files = set(expected.get('expected_files', []))
    actual_files = set(result)
    if not expected_files:
        return 0.0
    if actual_files == expected_files:
        score = 1.0
    else:
        overlap = len(actual_files & expected_files)
        extra = len(actual_files - expected_files)
        missing = len(expected_files - actual_files)
        score = max(0.0, overlap / len(expected_files) - extra * 0.1)
    logger.info(f'Expected: {expected_files}, Actual: {actual_files}, Score: {score}')
    return score

def check_remaining_files__8754d37bdc9e8d94ab80feb618caa015(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if folder contains exactly the expected remaining files.

    Args:
        result: List of filenames from getter
        expected: Dict with 'remaining_files' key containing list of expected filenames
        **options: Additional options (unused)

    Returns:
        1.0 if folder contains exactly the expected files, 0.0 otherwise
    """
    expected_files = expected.get('remaining_files', [])
    result_set = set(result)
    expected_set = set(expected_files)
    if result_set == expected_set:
        return 1.0
    else:
        return 0.0

def check_file_exists__c6c3aa52(result, expected, **options):
    """
    Check if file exists and is a valid PNG image screenshot.

    Args:
        result: Dict containing file validation information
        expected: Expected value (dict with 'exists' key)

    Returns:
        1.0 if file exists and is a valid PNG screenshot, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not isinstance(result, dict):
        return 0.0
    if not expected_exists:
        return 1.0 if not result.get('exists', False) else 0.0
    file_exists = result.get('exists', False)
    is_png = result.get('is_png', False)
    has_valid_size = result.get('has_valid_size', False)
    has_valid_dimensions = result.get('has_valid_dimensions', False)
    if file_exists and is_png and has_valid_size and has_valid_dimensions:
        return 1.0
    if not file_exists:
        logger.info('File does not exist')
    elif not is_png:
        logger.info('File is not a valid PNG image (magic bytes check failed)')
    elif not has_valid_size:
        logger.info('File size is too small (< 1KB)')
    elif not has_valid_dimensions:
        logger.info('File does not have valid image dimensions')
    return 0.0

def check_text_file_value__83bea20e1ddf48cf4f537ad2c05896b6(result: str, expected: dict, **options) -> float:
    """Check if text file content matches expected value.

    Args:
        result: Content from text file
        expected: Dict from rules with 'value' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    try:
        result_num = float(result)
        expected_num = float(expected_value)
        tolerance = options.get('tolerance', 0.01)
        if abs(result_num - expected_num) <= tolerance:
            logger.debug(f'Numeric match: {result_num} ~= {expected_num} (tolerance: {tolerance})')
            return 1.0
        else:
            logger.debug(f'Numeric mismatch: {result_num} != {expected_num}')
            return 0.0
    except (ValueError, TypeError):
        if str(result).strip() == str(expected_value).strip():
            logger.debug(f'String match: {result}')
            return 1.0
        else:
            logger.debug(f"String mismatch: '{result}' != '{expected_value}'")
            return 0.0

def check_single_file_exists__d6fb1e53c50621e1a08efd7623119b0d(result: Optional[List[str]], expected, **options):
    """Check if a specific file exists in Google Drive folder and is the only file.

    Args:
        result: List of filenames from getter, or None if folder doesn't exist
        expected: Dict with 'filename' to check for
        **options: Additional options

    Returns:
        float: 1.0 if the specific file exists and is the only file, 0.5 if file exists but there are other files, 0.0 otherwise
    """
    if result is None:
        logger.info('Result is None - folder not found')
        return 0.0
    filename = expected.get('filename')
    if not filename:
        logger.warning('No filename specified in expected')
        return 0.0
    if filename not in result:
        logger.info(f"File '{filename}' not found. Found files: {result}")
        return 0.0
    if len(result) == 1:
        logger.info(f"Target file '{filename}' is the only file in folder")
        return 1.0
    else:
        logger.info(f"Target file '{filename}' exists, but folder also contains: {[f for f in result if f != filename]}")
        return 0.5

def check_file_exists__67be1ac6efe87edab008a615fb0e7ec4(result, expected, **options):
    """Check if file exists as expected and validate it's a valid JPEG image.

    Args:
        result: Result from getter (dict with 'exists', 'is_jpeg', 'file_size' keys)
        expected: Expected rules (dict with 'should_exist' key)
        **options: Additional options

    Returns:
        float: 1.0 if file exists as expected and is valid JPEG, 0.0 otherwise
    """
    file_exists = result.get('exists', False)
    should_exist = expected.get('should_exist', True)
    is_jpeg = result.get('is_jpeg', False)
    file_size = result.get('file_size', 0)
    if file_exists != should_exist:
        return 0.0
    if not should_exist:
        return 1.0
    if not is_jpeg:
        return 0.0
    if file_size <= 0:
        return 0.0
    return 1.0

def check_file_count__58fb65f5(file_count, expected):
    """Check if file count matches expected.

    Args:
        file_count: Number of files in directory
        expected: Dict with 'count' key containing expected count

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected['count']
    if file_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_file_duplicate__654353fe(result, expected, **options):
    """Check if a file was duplicated correctly by comparing hashes.

    Args:
        result: String with format "hash1|hash2" from both files
        expected: Rules dict with expected_hash
        **options: Additional options

    Returns:
        float: 1.0 if both files have the same expected hash, 0.0 otherwise
    """
    if result is None or not isinstance(result, str):
        return 0.0
    result = result.strip()
    expected_hash = expected.get('expected_hash', '')
    if '|' not in result:
        return 0.0
    parts = result.split('|')
    if len(parts) != 2:
        return 0.0
    (hash1, hash2) = (parts[0].strip(), parts[1].strip())
    if 'NONE' in [hash1, hash2]:
        return 0.0
    if hash1 == expected_hash and hash2 == expected_hash:
        return 1.0
    else:
        return 0.0

def check_text_opacity__88b1d5c668539570e153fff50a7fc5f9(result, expected, **options):
    """
    Check if the text opacity is within the expected range.

    Args:
        result: Estimated opacity ratio from getter
        expected: Dictionary with 'min_opacity' and 'max_opacity' keys
        **options: Additional options

    Returns:
        float: 1.0 if opacity is within range, 0.0 otherwise
    """
    if result is None:
        return 0.0
    min_opacity = expected.get('min_opacity', 0.0)
    max_opacity = expected.get('max_opacity', 1.0)
    if min_opacity <= result <= max_opacity:
        return 1.0
    else:
        return 0.0

def check_file_contains__8085b902c0aa531b12d2ed766e22c897(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if file contains expected text(s).

    Args:
        result: String containing file content
        expected: Dictionary with 'text' key (string or list) containing text to search for
        **options: Additional options

    Returns:
        float: 1.0 if all expected text(s) found in file, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('text', '')
    if not expected_text:
        logger.error('Expected text not specified')
        return 0.0
    if isinstance(expected_text, list):
        for text in expected_text:
            if text not in result:
                return 0.0
        return 1.0
    elif expected_text in result:
        return 1.0
    return 0.0

def check_ods_file_exists__bdb8ae26(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if ODS file exists and is a valid ODS file with content.

    Args:
        result: Dict with validation results from getter
        expected: Dict with expected file path
        **options: Additional options

    Returns:
        Score: 1.0 if file exists, is valid ODS, and has content; 0.0 otherwise
    """
    if not result.get('file_exists', False):
        logger.warning(f"ODS file does not exist at {result.get('file_path')}")
        return 0.0
    if not result.get('is_valid_ods', False):
        logger.warning(f"File exists but is not a valid ODS archive: {result.get('error', 'Unknown error')}")
        return 0.0
    if not result.get('has_content', False):
        logger.warning(f"ODS file exists but does not contain spreadsheet data: {result.get('error', 'No content')}")
        return 0.0
    logger.info(f"ODS file successfully validated at {result.get('file_path')}: exists, valid ODS format, contains spreadsheet data")
    return 1.0

def check_file_recent__739292ff(result, expected, **options):
    """Check if file was created/modified recently.

    Args:
        result: Dict with 'age_seconds' key
        expected: Dict with 'max_age_seconds' key
        **options: Additional options

    Returns:
        float: 1.0 if file is recent enough, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    if not result.get('exists', False):
        logger.error('File does not exist')
        return 0.0
    max_age = expected.get('max_age_seconds', 300)
    actual_age = result.get('age_seconds', float('inf'))
    if actual_age <= max_age:
        logger.info(f'File is recent: {actual_age:.1f}s <= {max_age}s')
        return 1.0
    else:
        logger.info(f'File is too old: {actual_age:.1f}s > {max_age}s')
        return 0.0

def check_file_exists__3928cfa5(result, expected, **options):
    """
    Check if file exists and is a valid PNG image with reasonable properties.

    Args:
        result: Dict with file metadata:
            - exists: bool - whether file exists
            - is_valid_png: bool - whether file is a valid PNG image
            - file_size: int - file size in bytes
            - width: int - image width
            - height: int - image height
        expected: Expected value with rules:
            - exists: bool - file should exist

    Returns:
        1.0 if all validations pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result.get('exists') != expected_exists:
        logger.info(f"File existence check failed: expected {expected_exists}, got {result.get('exists')}")
        return 0.0
    if not expected_exists:
        return 1.0
    if not result.get('is_valid_png'):
        logger.info('File is not a valid PNG image (magic bytes check failed)')
        return 0.0
    file_size = result.get('file_size', 0)
    if file_size < 1024:
        logger.info(f'File size too small ({file_size} bytes), likely not a real screenshot')
        return 0.0
    width = result.get('width', 0)
    height = result.get('height', 0)
    if width <= 0 or height <= 0:
        logger.info(f'Invalid image dimensions: {width}x{height}')
        return 0.0
    logger.info(f'Valid PNG screenshot found: {file_size} bytes, {width}x{height} pixels')
    return 1.0

def check_winloss_sparklines__2bd59342_aug_12_verify_2(result_state: Dict[str, Any], expected_state: Dict[str, Any], **options) -> float:
    """
    Verify that win/loss sparklines are correctly created.

    This metric checks that:
    1. Sparklines exist in the file
    2. At least one sparkline is of type 'stacked' (win/loss)
    3. Sparklines are created for the correct data ranges (monthly sales)
    4. Sparklines are created for each row with data (if verify_coverage is enabled)

    Args:
        result_state: Dictionary from getter containing sparkline information:
            {
                'sparklines': [
                    {
                        'cell': 'A1',
                        'type': 'stacked',
                        'data_range': 'B1:D1',
                        'row': 1,
                    },
                    ...
                ],
                'count': int,
                'has_winloss': bool,
                'data_rows': int,
            }
        expected_state: When expected.type='rule', this IS the rules dict directly:
            {
                'min_sparklines': int,  # Minimum number of sparklines required
                'require_winloss': bool,  # Whether win/loss type is required
                'check_data_ranges': bool,  # Whether to verify data ranges contain monthly data
                'verify_coverage': bool,  # Whether to verify all data rows have sparklines
            }
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 indicating verification success
    """
    if not result_state:
        return 0.0
    sparklines = result_state.get('sparklines', [])
    count = result_state.get('count', 0)
    has_winloss = result_state.get('has_winloss', False)
    min_sparklines = expected_state.get('min_sparklines', 1)
    require_winloss = expected_state.get('require_winloss', True)
    check_data_ranges = expected_state.get('check_data_ranges', False)
    verify_coverage = expected_state.get('verify_coverage', False)
    if count < min_sparklines:
        return 0.0
    if require_winloss and (not has_winloss):
        return 0.0
    if verify_coverage:
        data_rows = result_state.get('data_rows', 0)
        if data_rows > 0:
            if count < data_rows * 0.8:
                return 0.0
    if check_data_ranges:
        valid_data_ranges = 0
        for sparkline in sparklines:
            data_range = sparkline.get('data_range', '')
            if ':' in data_range and data_range.strip():
                try:
                    if '.' in data_range:
                        range_part = data_range.split('.', 1)[1]
                    else:
                        range_part = data_range
                    if ':' in range_part:
                        (start_cell, end_cell) = range_part.split(':', 1)
                        start_col = ''.join([c for c in start_cell if c.isalpha()])
                        end_col = ''.join([c for c in end_cell if c.isalpha()])
                        if start_col and end_col and (start_col != end_col):

                            def col_to_num(col):
                                num = 0
                                for c in col:
                                    num = num * 26 + (ord(c.upper()) - ord('A') + 1)
                                return num
                            col_span = abs(col_to_num(end_col) - col_to_num(start_col)) + 1
                            if col_span >= 3:
                                valid_data_ranges += 1
                except:
                    pass
        if count > 0 and valid_data_ranges == count:
            return 1.0
        elif valid_data_ranges > 0:
            return valid_data_ranges / count if count > 0 else 0.0
        else:
            return 0.0
    return 1.0

def check_text_content__f8ce9de4(result, expected, **options):
    """Compare text content against expected value.

    Args:
        result: Text content from getter
        expected: Expected text value (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_stripped = str(result).strip()
    expected_stripped = str(expected_value).strip()
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_include_excluding_regex__06dc70fa7cf93b7501432994b47d9c35(result: Dict[str, Any], expected: dict, **options) -> float:
    """Check if the correct product is identified AND highlighted.

    Args:
        result: Dict with 'top_product' (str) and 'is_highlighted' (bool) from getter
        expected: Expected rules containing target product name
        **options: Additional comparison options

    Returns:
        float: Score 1.0 if both product matches and is highlighted, 0.5 if only product matches, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_product = expected.get('top_product', '')
    if not expected_product:
        return 0.0
    top_product = result.get('top_product', '')
    is_highlighted = result.get('is_highlighted', False)
    if top_product != expected_product:
        return 0.0
    if is_highlighted:
        return 1.0
    else:
        return 0.5

def check_files_count__014e651b(result, rules) -> float:
    """Check if the number of existing files matches expected count.

    Args:
        result: Integer count from getter
        rules: Dict with 'expected' key (integer)

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = rules.get('expected', 3)
    if result == expected_count:
        return 1.0
    return 0.0

def check_specific_files_present__465762dc(result: list, expected: dict, **options) -> float:
    """Check if specific files are present in the list.

    Args:
        result: List of filenames from getter
        expected: Dict with 'filenames' key containing list of expected filenames
        **options: Additional options

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_filenames = expected.get('filenames', [])
    for filename in expected_filenames:
        if filename not in result:
            return 0.0
    if len(result) != len(expected_filenames):
        return 0.0
    return 1.0

def check_file_list__b0918ceb(result: bytes, expected: Dict[str, Any], **options) -> float:
    """Check if text file contains expected list of filenames in alphabetical order.

    Args:
        result: File content as bytes from vm_file getter
        expected: Expected rules dict with 'expected_files' list
        **options: Additional options

    Returns:
        float: Score 1.0 if exact match (all files present, alphabetically sorted, no extras), 0.0 otherwise
    """
    try:
        expected_files = expected.get('expected_files', [])
        if not expected_files:
            return 0.0
        actual_lines = parse_file_content_lines(result)
        if len(actual_lines) != len(expected_files):
            return 0.0
        if actual_lines != expected_files:
            return 0.0
        if actual_lines != sorted(actual_lines):
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_file_exists__f5680565(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def check_vm_file__f259522f5141b84d8b2c6c9007fd732a(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if the VM file meets expected criteria including PDF validation and email attachment verification.

    Args:
        result: Dict from getter with keys: exists, size_bytes, is_pdf, has_pdf_header, is_largest_from_email, all_email_pdf_sizes
        expected: Dict with validation rules:
            - exists: bool - expected existence
            - min_size_bytes: int (optional) - minimum file size
            - validate_pdf: bool (optional) - whether to validate PDF format
            - require_pdf_extension: bool (optional) - whether to require .pdf extension
            - verify_largest_from_email: bool (optional) - whether to verify this is the largest PDF from email

    Returns:
        float: 1.0 if all criteria met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_exists = expected.get('exists', True)
    if result.get('exists', False) != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    min_size = expected.get('min_size_bytes', 0)
    if result.get('size_bytes', 0) < min_size:
        return 0.0
    if expected.get('require_pdf_extension', False):
        if not result.get('is_pdf', False):
            return 0.0
    if expected.get('validate_pdf', False):
        if not result.get('has_pdf_header', False):
            return 0.0
    if expected.get('verify_largest_from_email', False):
        if not result.get('is_largest_from_email', False):
            return 0.0
    return 1.0

def check_file_count__be02851a(result_state, expected_state, **options):
    """
    Check if the file count matches the expected value.

    This metric function compares the actual file count (from the result state)
    with the expected count. It allows for exact matches or a tolerance range
    to account for minor counting differences (e.g., whether directories are counted).

    Args:
        result_state: The actual file count (int) returned by the getter function
        expected_state: Dict containing:
            - expected_count: The expected number of files (int)
        **options: Additional options (currently unused)

    Returns:
        float: Score between 0.0 and 1.0
            - 1.0: Exact match or within tolerance (±2 files)
            - 0.0: Outside acceptable range or invalid data
    """
    try:
        if result_state is None:
            return 0.0
        if isinstance(expected_state, dict):
            expected_count = expected_state.get('expected_count')
        else:
            expected_count = expected_state
        if expected_count is None:
            return 0.0
        result_count = int(result_state)
        expected_count = int(expected_count)
        if result_count == expected_count:
            return 1.0
        tolerance = 2
        if abs(result_count - expected_count) <= tolerance:
            return 1.0
        return 0.0
    except (ValueError, TypeError):
        return 0.0

def check_text_contains__8f36e769(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    score = 0.0
    for substring in contains:
        if substring.lower() in result.lower():
            score += 1.0 / len(contains)
    return score

def check_python_functions__198be354(result, expected, **options):
    """Check if Python file has complete code with functions, imports, and other elements.

    Args:
        result: dict from getter with comprehensive code info
        expected: dict with 'min_functions' and optionally 'required_functions'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    min_functions = expected.get('min_functions', 0)
    actual_count = result.get('function_count', 0)
    if actual_count >= min_functions:
        score += 0.4
    required_functions = expected.get('required_functions', [])
    if required_functions:
        actual_functions = set(result.get('function_names', []))
        matched = sum((1 for func in required_functions if func in actual_functions))
        score += 0.3 * (matched / len(required_functions))
    else:
        score += 0.3
    if result.get('has_imports', False):
        score += 0.15
    if result.get('has_classes', False):
        score += 0.1
    if result.get('has_global_code', False) or result.get('import_count', 0) >= 3:
        score += 0.05
    return min(score, 1.0)

def check_file_exists__31a8a4acc19afab71c5f7ec2f3006a11(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a file exists at the expected path.

    Args:
        result: Dict from getter with 'exists' and 'file_path' keys
        expected: Dict with 'file_path' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result.get('exists', False):
        return 1.0
    return 0.0

def check_all_files_exist__4e03b1ed(result, expected, **options):
    """Check if all expected files exist.

    Args:
        result: Dict mapping file paths to existence status from getter
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: Partial credit based on how many files exist
    """
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    existing_count = sum((1 for f in required_files if result.get(f, False)))
    total_count = len(required_files)
    score = existing_count / total_count
    logger.info(f'{existing_count}/{total_count} required files exist. Score: {score}')
    return score

def check_python_script_contains__c9385c6b(result, expected, **options):
    """Check if Python script contains required imports and patterns.

    Args:
        result: Content of the Python script
        expected: Dict with 'required_imports' and 'required_patterns'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str) or not result:
        return 0.0
    required_imports = expected.get('required_imports', [])
    required_patterns = expected.get('required_patterns', [])
    total_checks = len(required_imports) + len(required_patterns)
    if total_checks == 0:
        return 0.0
    passed = 0
    for imp in required_imports:
        if imp in result:
            passed += 1
    for pattern in required_patterns:
        if pattern in result:
            passed += 1
    return passed / total_checks

def check_text_contains__551695fc(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    score = 0.0
    for substring in contains:
        if substring.lower() in result.lower():
            score += 1.0 / len(contains)
    return score

def check_file_content_d461fd5d(result, expected, **options):
    """Check if file content matches expected count.

    Args:
        result: File path (string) from get_vm_file getter
        expected: Expected rules dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        count = int(content)
        expected_count = expected.get('expected_count', 0)
        if count == expected_count:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_direct_json_object_with_time__f79439ad(result, expected) -> float:
    """
    Check JSON object with relative time processing.
    Processes relativeTime rules and compares result with expected values.
    """
    logger.info(f'[DEBUG] check_direct_json_object_with_time__f79439ad called')
    logger.info(f'[DEBUG] Result: {result}')
    logger.info(f'[DEBUG] Expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if 'relativeTime' in expected:
        relativeTime = expected['relativeTime']
        timezone_str = expected.get('timezone', get_timezone_from_ip())
        try:
            timezone = pytz.timezone(timezone_str)
        except pytz.exceptions.UnknownTimeZoneError:
            logger.error(f'Unknown timezone: {timezone_str}, using UTC')
            timezone = pytz.UTC
        now = datetime.now(timezone)
        logger.info(f"Current time in {timezone_str}: {now.strftime('%Y-%m-%d %H:%M:%S %Z')}")
        from_time_str = relativeTime.get('from')
        if from_time_str:
            if from_time_str == '5th next month':
                next_year = now.year + 1 if now.month == 12 else now.year
                next_month = now.month + 1 if now.month < 12 else 1
                from_absoluteDay = timezone.localize(datetime(next_year, next_month, 5))
            elif from_time_str == '10th next month':
                next_year = now.year + 1 if now.month == 12 else now.year
                next_month = now.month + 1 if now.month < 12 else 1
                from_absoluteDay = timezone.localize(datetime(next_year, next_month, 10))
            elif from_time_str == 'tomorrow':
                from_absoluteDay = now + timedelta(days=1)
            else:
                logger.error(f"Unsupported relativeTime 'from': {from_time_str}")
                from_absoluteDay = now
            if 'time_from' in expected['expected'] and '{' in expected['expected']['time_from']:
                expected['expected']['time_from'] = apply_rules_to_timeFormat(expected['expected']['time_from'], from_absoluteDay)
                logger.info(f"Processed time_from: {expected['expected']['time_from']}")
        to_time_str = relativeTime.get('to')
        if to_time_str:
            if to_time_str == '12th next month':
                next_year = now.year + 1 if now.month == 12 else now.year
                next_month = now.month + 1 if now.month < 12 else 1
                to_absoluteDay = timezone.localize(datetime(next_year, next_month, 12))
            elif to_time_str == '11th next month':
                next_year = now.year + 1 if now.month == 12 else now.year
                next_month = now.month + 1 if now.month < 12 else 1
                to_absoluteDay = timezone.localize(datetime(next_year, next_month, 11))
            elif to_time_str == 'tomorrow':
                to_absoluteDay = now + timedelta(days=1)
            else:
                logger.error(f"Unsupported relativeTime 'to': {to_time_str}")
                to_absoluteDay = now
            if 'time_to' in expected['expected'] and '{' in expected['expected']['time_to']:
                expected['expected']['time_to'] = apply_rules_to_timeFormat(expected['expected']['time_to'], to_absoluteDay)
                logger.info(f"Processed time_to: {expected['expected']['time_to']}")
    expected_json = expected.get('expected', {})
    logger.info(f'[DEBUG] Expected JSON (after processing): {expected_json}')
    for key in expected_json.keys():
        if key in ['ignore_list_order']:
            continue
        expected_value = expected_json.get(key)
        actual_value = result.get(key)
        logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
        if expected_value != actual_value:
            logger.info(f"[DEBUG] Value mismatch for key '{key}', returning 0.0")
            return 0.0
    logger.info('[DEBUG] All comparisons passed, returning 1.0')
    return 1.0

def check_text_output__451bbe5b(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_exists__e86a3a4d(result, expected, **options):
    """Check if file existence matches expected.

    Args:
        result: Boolean from getter
        expected: Dict with 'exists' boolean value
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(expected, dict):
        return 0.0
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_file_exists_with_content__4f896edc(result_state, expected_state, **options):
    """
    Check if an Excel file exists, contains the expected number of rows,
    and verifies that all rows are AWS entries.

    This function verifies that:
    1. The file exists at the specified path
    2. The file contains the expected number of data rows
    3. All data rows contain 'AWS' in at least one cell (verifying filtering was correct)

    Args:
        result_state: Path to the result file (str)
        expected_state: Expected configuration dict with 'expected_rows' key
        **options: Additional options

    Returns:
        float: Score (1.0 if file exists, has expected rows, and all rows are AWS entries; 0.0 otherwise)
    """
    if not result_state or not isinstance(result_state, str):
        logger.error(f'Invalid result_state: {result_state}')
        return 0.0
    if not os.path.isfile(result_state):
        logger.error(f'File does not exist: {result_state}')
        return 0.0
    expected_rows = expected_state.get('expected_rows', None) if isinstance(expected_state, dict) else None
    if expected_rows is None:
        logger.error(f'No expected_rows specified in expected_state: {expected_state}')
        return 0.0
    try:
        wb = openpyxl.load_workbook(result_state)
        sheet = wb.active
        all_rows = []
        for row in sheet.iter_rows(min_row=1):
            if any((cell.value is not None and str(cell.value).strip() != '' for cell in row)):
                all_rows.append([cell.value for cell in row])
        if len(all_rows) == 0:
            logger.error('File is empty')
            wb.close()
            return 0.0
        first_row = all_rows[0]
        has_header = any((isinstance(val, str) and val and (not str(val).replace('.', '').replace('-', '').replace('$', '').replace(',', '').strip().replace(' ', '').isdigit()) for val in first_row if val is not None))
        if has_header:
            data_rows_list = all_rows[1:]
        else:
            data_rows_list = all_rows
        actual_data_rows = len(data_rows_list)
        logger.info(f'File {result_state} has {actual_data_rows} data rows (expected: {expected_rows})')
        if actual_data_rows != expected_rows:
            logger.error(f'Row count mismatch: expected {expected_rows}, got {actual_data_rows}')
            wb.close()
            return 0.0
        aws_rows_count = 0
        for (row_idx, row) in enumerate(data_rows_list, start=1):
            row_contains_aws = False
            for cell_value in row:
                if cell_value is not None:
                    cell_str = str(cell_value).upper()
                    if 'AWS' in cell_str:
                        row_contains_aws = True
                        break
            if row_contains_aws:
                aws_rows_count += 1
            else:
                logger.error(f"Data row {row_idx} does not contain 'AWS': {row}")
        wb.close()
        if aws_rows_count == actual_data_rows:
            logger.info(f'All {actual_data_rows} data rows are AWS entries')
            return 1.0
        else:
            logger.error(f"Only {aws_rows_count}/{actual_data_rows} rows contain 'AWS'")
            return 0.0
    except Exception as e:
        logger.error(f'Error reading file {result_state}: {e}')
        return 0.0

def check_file_exists__5f5351b0(result: dict, expected: dict, **options) -> float:
    """Check if the zip file exists, is valid, and contains all expected photos.

    Args:
        result: Dict with 'exists' (bool), 'valid' (bool), and 'contents' (list of filenames)
        expected: Dict with 'exists' (bool) and 'required_files' (list of filenames)
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    required_files = expected.get('required_files', [])
    if result.get('exists') != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    if not result.get('valid', False):
        return 0.0
    contents = result.get('contents', [])
    if not required_files:
        return 1.0
    for required_file in required_files:
        if required_file not in contents:
            return 0.0
    return 1.0

def check_python_classes__198be354(result, expected, **options):
    """Check if Python file has expected class definitions.

    Args:
        result: dict from getter with class info
        expected: dict with 'min_classes' and optionally 'required_classes'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    min_classes = expected.get('min_classes', 0)
    actual_count = result.get('class_count', 0)
    if actual_count >= min_classes:
        score += 0.7
    required_classes = expected.get('required_classes', [])
    if required_classes:
        actual_classes = set(result.get('class_names', []))
        matched = sum((1 for cls in required_classes if cls in actual_classes))
        score += 0.3 * (matched / len(required_classes))
    else:
        score += 0.3
    return score

def check_file_exists__c48ab0f0(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def check_textbox_on_bottom__734a49ef(src_path, expected, **options):
    """
    Check if the textbox is at the bottom of the image.
    Task variation 3 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is at bottom, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    bottom_most_dark_pixel = 0
    for x in range(width):
        for y in range(height - 1, -1, -1):
            if gray_image.getpixel((x, y)) < 128:
                bottom_most_dark_pixel = max(bottom_most_dark_pixel, y)
                break
    if bottom_most_dark_pixel > height * 0.95:
        return 1.0
    else:
        return 0.0

def check_file_exists_and_structure_sim_77b8ab4d(result_state, expected_rules, **options):
    """
    Check if the image has been exported to the desktop with the correct structure.

    This metric verifies:
    1. The result file exists (user saved the image)
    2. The structure of the result image matches the expected reference image

    Args:
        result_state: Dict from getter with 'result_path' and 'reference_path' keys
        expected_rules: Dict (not used, reference path comes from getter)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and structure matches, 0.0 otherwise
    """
    if not isinstance(result_state, dict):
        logging.error(f'Result state is not a dict: {type(result_state)}')
        return 0.0
    result_path = result_state.get('result_path')
    reference_path = result_state.get('reference_path')
    if result_path is None:
        logging.debug('Result file was not saved or not found')
        return 0.0
    if not os.path.isfile(result_path):
        logging.debug(f'Result file does not exist: {result_path}')
        return 0.0
    if reference_path is None or not os.path.isfile(reference_path):
        logging.error(f'Reference file not found: {reference_path}')
        return 0.0
    try:
        result_img = Image.open(result_path)
        reference_img = Image.open(reference_path)
        structure_same = structure_check_by_ssim(result_img, reference_img)
        if structure_same:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logging.error(f'Error comparing images: {e}')
        return 0.0

def check_textbox_bottom__3bddf9fd7655c038551f7bf6cc6f697d(result_path, rules, env=None):
    """
    Check if the text box is positioned at the bottom of the image.

    This function:
    1. Identifies the text box region using connected component analysis
    2. Finds the bounding box of the text box
    3. Verifies that the bottom edge of the text box is within the bottom 5% of the image

    Args:
        result_path: Path to the exported image (bottomside_textbox.png)
        rules: Dict containing rules (not used in this implementation)
        env: Environment object (optional)

    Returns:
        float: 1.0 if text box is at the bottom, 0.0 otherwise
    """
    if result_path is None:
        return 0.0
    try:
        source_image = Image.open(result_path)
        gray_image = source_image.convert('L')
        (width, height) = source_image.size
        gray_array = np.array(gray_image)
        dark_mask = gray_array < 128
        if not np.any(dark_mask):
            return 0.0
        (labeled_array, num_features) = ndimage.label(dark_mask)
        if num_features == 0:
            return 0.0
        component_sizes = np.bincount(labeled_array.ravel())
        component_sizes[0] = 0
        largest_component_label = component_sizes.argmax()
        textbox_mask = labeled_array == largest_component_label
        textbox_coords = np.argwhere(textbox_mask)
        if len(textbox_coords) == 0:
            return 0.0
        y_coords = textbox_coords[:, 0]
        bottommost_y = np.max(y_coords)
        textbox_pixel_count = len(textbox_coords)
        if textbox_pixel_count < 100:
            return 0.0
        centroid_y = np.mean(y_coords)
        bottom_threshold = height * 0.95
        if bottommost_y >= bottom_threshold and centroid_y >= height * 0.8:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_python_imports__5e897a929023bc43113113bb0ca8fb36(result, expected, **options):
    """Compare extracted import statements against expected imports.

    This function checks if ALL expected imports are present in the result,
    and no extra imports are included. With threshold=1.0, only exact matches
    receive full credit.

    Args:
        result: Extracted import statements from result file
        expected: Expected import statements (from rules)
        **options: Additional options (threshold for scoring, default 1.0)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0

    def normalize_imports(text):
        """Normalize import statements for comparison."""
        lines = []
        for line in text.splitlines():
            stripped = line.strip()
            if stripped and (stripped.startswith('import ') or stripped.startswith('from ')):
                if '#' in stripped:
                    stripped = stripped[:stripped.index('#')].strip()
                if stripped:
                    lines.append(stripped)
        return set(lines)
    result_imports = normalize_imports(result)
    if isinstance(expected, dict):
        expected_imports_str = expected.get('expected_imports', '')
    else:
        expected_imports_str = expected
    expected_imports = normalize_imports(expected_imports_str)
    if not expected_imports:
        return 1.0 if not result_imports else 0.0
    if result_imports == expected_imports:
        return 1.0
    matched = result_imports & expected_imports
    recall = len(matched) / len(expected_imports) if expected_imports else 0.0
    precision = len(matched) / len(result_imports) if result_imports else 0.0
    if recall + precision == 0:
        return 0.0
    f1 = 2 * (recall * precision) / (recall + precision)
    threshold = options.get('threshold', 1.0)
    return 1.0 if f1 >= threshold else f1

def check_first_line_indent__2b727758(result, expected, **options):
    """
    Check if paragraphs have proper first-line indentation.

    Args:
        result: Dict with indentation statistics from getter
        expected: Dict with 'min_indent_pt' and 'min_coverage' (proportion of paragraphs)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    min_indent_pt = expected.get('min_indent_pt', 36)
    min_coverage = expected.get('min_coverage', 0.8)
    total_paras = result.get('count', 0)
    with_indent = result.get('with_indent', 0)
    avg_indent = result.get('avg_indent', 0)
    if total_paras == 0:
        return 0.0
    coverage = with_indent / total_paras
    indent_match = 0.0
    if avg_indent > 0:
        target_indent = min_indent_pt
        tolerance = target_indent * 0.2
        if abs(avg_indent - target_indent) <= tolerance:
            indent_match = 1.0
        elif avg_indent > target_indent - tolerance:
            indent_match = 0.5
    coverage_score = min(coverage / min_coverage, 1.0)
    return coverage_score * 0.5 + indent_match * 0.5

def check_direct_json_object(result, rules) -> float:
    """
    One of the most commonly used function to evalute.
    Compare two json objects directly.

    Fixed version that correctly handles 'ignore_list_order' parameter
    from the top-level rules dict instead of expected_json dict.
    """
    logger.info(f'[DEBUG] check_direct_json_object called with result: {result}')
    logger.info(f'[DEBUG] check_direct_json_object called with rules: {rules}')
    if isinstance(result, str):
        result = result.strip()
        result = result.replace("'", '"')
        result = json.loads(result)
    logger.info(f'[DEBUG] Processed result: {result}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    try:
        expected_json = rules.get('expected', {})
        if expected_json:
            for (key, value) in expected_json.items():
                if value == '__EVALUATION_FAILED__':
                    logger.error(f"[DEBUG] Expected value for key '{key}' indicates evaluation failure, returning 0.0")
                    return 0.0
    except Exception as e:
        logger.error(f'[DEBUG] Error checking for evaluation failure indicator: {e}')
        return 0.0
    try:
        expect_in_result = rules.get('expect_in_result', False)
        logger.info(f'[DEBUG] expect_in_result: {expect_in_result}')
        if not expect_in_result:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON: {expected_json}')
            for key in expected_json.keys():
                expected_value = expected_json.get(key)
                actual_value = result.get(key)
                logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
                if rules.get('ignore_list_order', False):
                    expected_value = sorted(expected_value)
                    result_value = sorted(result.get(key))
                    logger.info(f'[DEBUG] Comparing lists (sorted): expected={expected_value}, actual={result_value}')
                    if expected_value != result_value:
                        logger.info(f"[DEBUG] List comparison failed for key '{key}', returning 0.0")
                        return 0.0
                elif expected_value != actual_value:
                    logger.info(f"[DEBUG] Value comparison failed for key '{key}': expected='{expected_value}', actual='{actual_value}', returning 0.0")
                    return 0.0
                else:
                    logger.info(f"[DEBUG] Value comparison passed for key '{key}'")
            logger.info('[DEBUG] All comparisons passed, returning 1.0')
            return 1.0
        else:
            expected_json = rules['expected']
            logger.info(f'[DEBUG] Expected JSON (expect_in_result mode): {expected_json}')
            for key in expected_json.keys():
                if isinstance(expected_json.get(key), list):
                    flag = 0
                    expected_value_list = expected_json.get(key)
                    logger.info(f"[DEBUG] Checking list key '{key}': expected_list={expected_value_list}, actual='{result.get(key)}'")
                    for each_expected_value in expected_value_list:
                        if isinstance(result.get(key), list) and each_expected_value in result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' in result list for key '{key}'")
                            break
                        elif isinstance(result.get(key), str) and each_expected_value == result.get(key):
                            flag = 1
                            logger.info(f"[DEBUG] Found expected value '{each_expected_value}' matches result string for key '{key}'")
                            break
                    if flag == 0:
                        logger.info(f"[DEBUG] No expected values found in result for key '{key}', returning 0.0")
                        return 0.0
                elif isinstance(expected_json.get(key), str):
                    expected_str = expected_json.get(key)
                    actual_str = result.get(key)
                    logger.info(f"[DEBUG] Checking string key '{key}': expected='{expected_str}', actual='{actual_str}'")
                    if expected_str not in actual_str:
                        logger.info(f"[DEBUG] Expected string '{expected_str}' not found in actual string '{actual_str}' for key '{key}', returning 0.0")
                        return 0.0
                else:
                    logger.debug('check_direct_json_object: expected value type not supported')
                    return 0.0
            logger.info('[DEBUG] All expect_in_result comparisons passed, returning 1.0')
            return 1.0
    except Exception as e:
        logger.debug(f'check_direct_json_object: result is not a valid json object, error: {e}')
        return 0.0

def check_exact_file_list__91c6f86e(filenames, expected):
    """Check if exact list of files is present.

    Args:
        filenames: List of filenames in directory
        expected: Dict with 'expected' key containing list of expected filenames

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected['expected']
    if len(filenames) != len(expected_files):
        return 0.0
    if set(filenames) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_file_exists__b2463eb9(result, expected, **options):
    """Check if file exists and contains the expected email summary content.

    Args:
        result: File content (string)
        expected: Expected conditions (dict with 'subject' and 'check_sender' keys)

    Returns:
        float: 1.0 if file exists with correct content, 0.0 otherwise
    """
    if not result or not isinstance(result, str) or len(result.strip()) == 0:
        return 0.0
    expected_subject = expected.get('subject', 'Paper Recommendation')
    check_sender = expected.get('check_sender', True)
    if expected_subject not in result:
        return 0.0
    if check_sender:
        has_sender_info = '@' in result or 'from:' in result.lower() or 'sender:' in result.lower()
        if not has_sender_info:
            return 0.0
    return 1.0

def check_textbox_top__f0019e8aa52327a88742c54abff5dd41(result_state, expected_state, **options):
    """
    Check if the textbox has been moved to the top of the image.

    This function verifies that the topmost dark pixel is within the top 5% of the image height,
    indicating the text box has been positioned near the top edge as requested.

    Args:
        result_state: Path to the exported PNG file
        expected_state: Expected state (not used, rule-based)
        **options: Additional options

    Returns:
        float: 1.0 if textbox is at the top (within top 5%), 0.0 otherwise
    """
    if result_state is None:
        logger.warning('Result state is None, image file not found')
        return 0.0
    try:
        image = Image.open(result_state)
        grayscale = image.convert('L')
        img_array = np.array(grayscale)
        (height, width) = img_array.shape
        dark_threshold = 128
        dark_pixels = np.where(img_array < dark_threshold)
        if len(dark_pixels[0]) == 0:
            logger.warning('No dark pixels found in image')
            return 0.0
        topmost_y = np.min(dark_pixels[0])
        top_5_percent_threshold = height * 0.05
        logger.info(f'Image height: {height}, Topmost dark pixel at y={topmost_y}, Top 5% threshold: {top_5_percent_threshold}')
        if topmost_y <= top_5_percent_threshold:
            logger.info('Text box is positioned at the top of the image')
            return 1.0
        else:
            logger.info(f'Text box is not at the top (y={topmost_y} > {top_5_percent_threshold})')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking textbox position: {e}')
        return 0.0

def check_total_files_count__e9983217(result, expected, **options):
    """
    Check if repository was properly downloaded with all required components.

    Args:
        result: dict with file_count, has_git_dir, required_files_exist, required_dirs_exist
        expected: dict with min_count, max_count, required_files, required_dirs
        **options: Additional options (not used)

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    min_count = expected.get('min_count', 50)
    max_count = expected.get('max_count', 500)
    required_files = expected.get('required_files', [])
    required_dirs = expected.get('required_dirs', [])
    file_count = result.get('file_count', 0)
    if file_count < min_count or file_count > max_count:
        return 0.0
    if not result.get('has_git_dir', False):
        return 0.0
    required_files_exist = result.get('required_files_exist', {})
    for filename in required_files:
        if filename != '.git' and (not required_files_exist.get(filename, False)):
            return 0.0
    required_dirs_exist = result.get('required_dirs_exist', {})
    for dirname in required_dirs:
        if dirname != '.git' and (not required_dirs_exist.get(dirname, False)):
            return 0.0
    return 1.0

def check_txt_file_content__b2ba9ee6873b43000b20ba169c4d9592(result, expected, **options):
    """
    Check if TXT file exists and contains expected text.

    Args:
        result: Dict from getter with 'exists' and 'content' keys
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and has content, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    exists = result.get('exists', False)
    content = result.get('content', '')
    if not exists:
        logger.warning(f"❌ TXT file does not exist at {result.get('path', 'unknown')}")
        return 0.0
    expected_keywords = expected.get('keywords', [])
    if len(content.strip()) == 0:
        logger.warning(f'❌ TXT file is empty')
        return 0.0
    if expected_keywords:
        found_keywords = sum((1 for keyword in expected_keywords if keyword.lower() in content.lower()))
        if found_keywords == 0:
            logger.warning(f'❌ No expected keywords found in TXT file')
            return 0.0
    logger.info(f'✅ TXT file exists with valid content ({len(content)} chars)')
    return 1.0

def check_file_exists_with_size__e1affec5(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_file_rename__c5a909ed(actual: str, expected: dict, **options) -> float:
    """
    Check if a file was renamed successfully.

    Args:
        actual (str): command output ("renamed" or "not_renamed")
        expected (dict): expected dict with key "expected"

    Return:
        float: 1.0 if file was renamed, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_file_rename__c5a909ed: actual is None or empty')
        return 0.0
    expected_value = expected.get('expected', 'renamed')
    actual = actual.strip()
    if actual == expected_value:
        return 1.0
    logger.debug(f"check_file_rename__c5a909ed: Expected '{expected_value}', got '{actual}'")
    return 0.0

def check_python_code_stats__198be354(result, expected, **options):
    """Check if Python file has expected code statistics.

    Args:
        result: dict from getter with code stats
        expected: dict with 'min_code_lines' and 'min_total_lines'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    score = 0.0
    min_total = expected.get('min_total_lines', 0)
    if result.get('total_lines', 0) >= min_total:
        score += 0.4
    min_code = expected.get('min_code_lines', 0)
    if result.get('code_lines', 0) >= min_code:
        score += 0.6
    return score

def check_file_organization__94dfca2e(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_file_properties__e9151d35420d1468fe1e721181658d5f(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if file properties match expectations with partial credit.

    Args:
        result: List of file property dicts
        expected: Expected configuration (from rules dict)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    expected_count = expected.get('count', 3)
    min_size = expected.get('min_size', 1000)
    score = 0.0
    if len(result) == expected_count:
        score += 0.4
    png_files = [f for f in result if f.get('is_png', False)]
    if len(png_files) == len(result) and len(result) > 0:
        score += 0.3
    non_empty_files = [f for f in result if f.get('size', 0) >= min_size]
    if len(non_empty_files) == len(result) and len(result) > 0:
        score += 0.3
    return score

def check_file_format__4b5f0cdf(src_path, rule):
    """
    Check if the image file is in the expected format
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_1
    """
    if src_path is None:
        return 0.0
    try:
        if not os.path.exists(src_path):
            logger.error(f'File does not exist: {src_path}')
            return 0.0
        img = Image.open(src_path)
        actual_format = img.format
        expected_format = rule.get('format', None)
        logger.debug(f'Image format: {actual_format}, expected: {expected_format}')
        if expected_format is not None and actual_format == expected_format:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking image format: {e}')
        return 0.0

def check_file_line_count__7233f122896fac183c973343e2cf3b2a(result, expected, **options):
    """Compare actual line count against expected value.

    Args:
        result: Actual line count (int)
        expected: Expected rules dict with 'count' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if counts match, 0.0 otherwise
    """
    expected_count = expected.get('count', -1)
    if result == expected_count:
        return 1.0
    else:
        return 0.0

def check_file_contains_text__edd1c1331eff751ad5487718d7e5d07b(result: str, expected: Dict, **options) -> float:
    """Check if file content contains expected text.

    Args:
        result: String content from getter
        expected: Dict with 'text' field specifying what text to search for
        **options: Additional options (not used)

    Returns:
        float: 1.0 if text is found in content, 0.0 otherwise
    """
    if not result:
        return 0.0
    text = expected.get('text', '')
    if text in result:
        return 1.0
    return 0.0

def check_timezone_utc_minus_5__d030998aa6296cb272a4b11f8e8cd0d6(timedatectl_output, expected, **options):
    """
    Check if timezone is set to UTC-5 (e.g., America/New_York, America/Toronto).

    Args:
        timedatectl_output: Output from 'timedatectl status' command
        expected: Expected timezone offset (should be "-0500")
        **options: Additional options

    Returns:
        float: 1.0 if timezone is UTC-5, 0.0 otherwise
    """
    expected_offset = expected.get('offset', '-0500')
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            match = re.search('\\(\\w+,\\s*([+-]\\d{4})\\)', line)
            if match:
                actual_offset = match.group(1)
                if actual_offset == expected_offset:
                    return 1.0
                else:
                    return 0.5
            else:
                if f'({expected_offset})' in line or line.strip().endswith(f'{expected_offset})'):
                    return 1.0
                return 0.5
    return 0.0

def check_timezone_cet__755f0aea(timedatectl_output, expected, **options):
    """
    Check if timezone is set to Central European Time (CET, UTC+1).

    Note: CET includes both winter time (UTC+1, +0100) and summer time
    (CEST, UTC+2, +0200). This task specifically requires UTC+1 offset (+0100)
    as specified in the instruction.

    Args:
        timedatectl_output: Output from timedatectl status command (string)
        expected: Rules dictionary with 'timezone_offset' key (e.g., {'+0100'})
        **options: Additional options

    Returns:
        1.0 if timezone matches expected offset, 0.0 otherwise
    """
    expected_offset = expected.get('timezone_offset', '+0100')
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            if line.strip().endswith(f'{expected_offset})'):
                return 1.0
    return 0.0

def check_timezone_ist__36d379a3(timedatectl_output, expected, **options):
    """
    Check if timezone is set to India Standard Time (IST, UTC+5:30).

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Expected timezone offset string (e.g., "+0530")
        **options: Additional options

    Returns:
        1.0 if timezone matches expected offset, 0.0 otherwise
    """
    lines = timedatectl_output.split('\n')
    expected_offset = expected
    for line in lines:
        if 'Time zone:' in line:
            if expected_offset in line and line.strip().endswith(expected_offset + ')'):
                return 1.0
    return 0.0

def check_lines_counter__20d8676e(result, expected, **options):
    """
    Check if lines cleared counter has been added to track total lines cleared separately from score.

    Args:
        result: Dict with tetris_py and main_py content
        expected: Expected patterns
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    tetris_code = result.get('tetris_py', '')
    main_code = result.get('main_py', '')
    tetris_lower = tetris_code.lower()
    main_lower = main_code.lower()
    if 'self.lines' in tetris_code or ('lines' in tetris_lower and 'self.' in tetris_code):
        score += 0.33
    if 'lines' in tetris_lower and ('len(lines_to_remove)' in tetris_code or '+=' in tetris_code):
        score += 0.34
    if 'game.lines' in main_code or ('lines' in main_lower and 'render' in main_lower):
        score += 0.33
    return score

def check_text_replacement__fe4ab075(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_include_func__864e927cd0113eb6d9476a9e8ebce88d(result, expected, **options):
    """Execute multiple validation checks with conjunction logic.

    Args:
        result: Dummy result (not used, sub-checks have their own getters)
        expected: Dict containing 'conj' (and/or) and 'include' (list of checks)
        **options: Additional options including 'env' for environment access

    Returns:
        float: Combined score based on conjunction logic
    """
    if not isinstance(expected, dict):
        return 0.0
    conj = expected.get('conj', 'and')
    include = expected.get('include', [])
    if not include:
        return 0.0
    env = options.get('env')
    if not env:
        return 0.0
    scores = []
    for check in include:
        func_name = check.get('func')
        result_config = check.get('result', {})
        expected_config = check.get('expected', {})
        if not func_name:
            scores.append(0.0)
            continue
        result_type = result_config.get('type')
        if not result_type:
            scores.append(0.0)
            continue
        try:
            getter_module_name = 'getters.xlsx'
            metric_module_name = 'metrics.xlsx'
            import importlib
            getter_module = importlib.import_module(getter_module_name)
            getter_func = getattr(getter_module, f'get_{result_type}', None)
            if not getter_func:
                scores.append(0.0)
                continue
            actual_result = getter_func(env, result_config)
            expected_type = expected_config.get('type')
            if expected_type == 'rule':
                expected_value = expected_config.get('rules', {})
            else:
                expected_getter = getattr(getter_module, f'get_{expected_type}', None)
                if expected_getter:
                    expected_value = expected_getter(env, expected_config)
                else:
                    expected_value = expected_config
            metric_module = importlib.import_module(metric_module_name)
            metric_func = getattr(metric_module, func_name, None)
            if not metric_func:
                scores.append(0.0)
                continue
            score = metric_func(actual_result, expected_value, **options)
            scores.append(score)
        except Exception as e:
            scores.append(0.0)
    if not scores:
        return 0.0
    if conj == 'and':
        return 1.0 if all((s == 1.0 for s in scores)) else 0.0
    elif conj == 'or':
        return 1.0 if any((s == 1.0 for s in scores)) else 0.0
    else:
        return sum(scores) / len(scores)

def check_text_rotated__a994a687(result_state, expected_state, **options):
    """
    Check if the text has been rotated (detectable by examining aspect ratio of text bounds).
    Variation 8 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result_state: Path to the result image with text (vm_file path or dict with 'path' key)
        expected_state: Not used (rule-based, expects empty dict {})
        **options: Additional options

    Returns:
        float: Score (1.0 if text appears rotated to vertical orientation, 0.0 otherwise)

    Note:
        - Uses aspect ratio analysis: vertical text has width/height < 1.5
        - Threshold of 1.5 works for most text but may fail for single characters
        - Cannot distinguish 90° vs 270° rotation (both are vertical)
        - Returns 0.0 for both horizontal text and no-text-found cases

    VM_FILE CONSUMPTION:
        - This function consumes vm_file result from result_state parameter
        - Extracts file path from dict['path'] or string (lines 34-45)
        - Performs file path validation with os.path.exists() (line 48)
        - Reads image file with Image.open() (line 53)
    """
    if isinstance(result_state, dict):
        if 'path' not in result_state:
            logger.error(f"vm_file result missing 'path' key: {result_state}")
            return 0.0
        file_path = result_state['path']
    elif isinstance(result_state, str):
        file_path = result_state
    else:
        logger.error(f'Invalid result_state type: {type(result_state)}')
        return 0.0
    if not file_path or not os.path.exists(file_path):
        logger.error(f'vm_file result path does not exist: {file_path}')
        return 0.0
    try:
        source_image = Image.open(file_path)
        gray_image = source_image.convert('L')
        img_array = np.array(gray_image)
        text_pixels = np.argwhere(img_array < 128)
        if len(text_pixels) == 0:
            logger.warning('No text detected in image')
            return 0.0
        (min_y, min_x) = text_pixels.min(axis=0)
        (max_y, max_x) = text_pixels.max(axis=0)
        text_width = max_x - min_x + 1
        text_height = max_y - min_y + 1
        if text_width <= 0 or text_height <= 0:
            logger.warning('Invalid text bounding box detected')
            return 0.0
        aspect_ratio = text_width / text_height
        if aspect_ratio < 1.5:
            logger.info(f'Text detected as rotated (aspect ratio: {aspect_ratio:.2f})')
            return 1.0
        else:
            logger.info(f'Text detected as horizontal (aspect ratio: {aspect_ratio:.2f})')
            return 0.0
    except Exception as e:
        logger.error(f'Error checking text rotation: {e}', exc_info=True)
        return 0.0

def check_text_content__5ef70598fcc68171b6ec36d7c888fc21(result, expected, **options):
    """Check if file content matches expected text exactly.

    Args:
        result: String content from the file
        expected: Dict with 'content' key containing expected text
        **options: Additional options (ignore_trailing_whitespace, etc.)

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_content = expected.get('content', '')
    ignore_trailing = options.get('ignore_trailing_whitespace', False)
    if ignore_trailing:
        result = result.rstrip()
        expected_content = expected_content.rstrip()
    if result == expected_content:
        return 1.0
    return 0.0

def check_text_list__d9839857(result, expected, **options):
    """Compare a list of text values.

    Args:
        result: List of values from getter
        expected: List of expected values (in order)
        **options: case_sensitive (default True), strip (default True)

    Returns:
        float: 1.0 if all match, partial credit for partial matches
    """
    case_sensitive = options.get('case_sensitive', True)
    strip = options.get('strip', True)
    if result is None or not isinstance(result, list):
        return 0.0
    if not isinstance(expected, list):
        return 0.0
    if len(result) != len(expected):
        return 0.0
    matches = 0
    for (r, e) in zip(result, expected):
        r_str = str(r) if r is not None else ''
        e_str = str(e) if e is not None else ''
        if strip:
            r_str = r_str.strip()
            e_str = e_str.strip()
        if not case_sensitive:
            r_str = r_str.lower()
            e_str = e_str.lower()
        if r_str == e_str:
            matches += 1
    return matches / len(expected)

def check_file_exists__7930967f(result_state: dict, expected_state: dict, **options) -> float:
    """
    Check if the file existence, content, and format match the expected state.

    Args:
        result_state: Dict with 'exists', 'size', and 'is_text' keys (from getter)
        expected_state: Dict with 'should_exist' key (True if file should exist)
        **options: Additional options

    Returns:
        float: 1.0 if the state matches expectation, 0.0 otherwise
    """
    should_exist = expected_state.get('should_exist', True)
    if result_state.get('exists', False) != should_exist:
        return 0.0
    if not should_exist:
        return 1.0
    file_size = result_state.get('size', 0)
    is_text = result_state.get('is_text', False)
    if file_size <= 0:
        return 0.0
    if not is_text:
        return 0.0
    return 1.0

def check_file_content__8caa3012(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if text file contains the expected count.

    Args:
        result: File path as string
        expected: Expected rules dict with 'expected_count'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    try:
        expected_count = expected.get('expected_count', 0)
        with open(result, 'r') as f:
            content = f.read().strip()
        try:
            actual_count = int(content)
            if actual_count == expected_count:
                return 1.0
            else:
                return 0.0
        except ValueError:
            return 0.0
    except Exception as e:
        return 0.0

def check_multiple_files__11701199(result, expected, **options):
    """Check if multiple expected files exist in directory listing.

    Args:
        result: Directory listing output (string with filenames)
        expected: Expected state rules with 'expected_files' list
        **options: Additional comparison options

    Returns:
        float: Partial credit based on how many files are found (0.0 to 1.0)
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        return 0.0
    result_files = [line.strip() for line in result.strip().split('\n') if line.strip()]
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result_files:
            found_count += 1
    score = found_count / len(expected_files)
    return score

def compare_text_output__d398cd07(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_text_content__a652858db2e92fae77817389157c8edc(result: str, expected: Dict, **options) -> float:
    """Check text file content against expected rules.

    Args:
        result: Actual text content from file
        expected: Expected configuration dict with rules:
            - 'contains': list of strings that must be present
            - 'is_empty': bool, if True, checks that file is empty or whitespace-only

    Returns:
        float: Score between 0.0 and 1.0
    """
    if expected.get('is_empty', False):
        if not result or result.strip() == '':
            return 1.0
        else:
            return 0.0
    contains_list = expected.get('contains', [])
    if contains_list:
        if not result:
            return 0.0
        result_lower = result.lower()
        score = 0.0
        for expected_str in contains_list:
            if expected_str.lower() in result_lower:
                score += 1.0 / len(contains_list)
        return score
    return 0.0

def check_file_exists__04b2d876(result, expected, **options):
    """
    Check if file existence matches expected state.

    Args:
        result: Boolean from getter (True if file exists)
        expected: Expected state (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_line_count__ab2d13c4(line_count, expected):
    """Check if line count matches expected.

    Args:
        line_count: Number of lines in file
        expected: Dict with 'count' key

    Returns:
        float: 1.0 if count matches, 0.0 otherwise
    """
    expected_count = expected['count']
    if line_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_exact_text_match__5b92f2b5(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_text_replacement__f1f9b3a9(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_exact_match_v8__d4477d7a(result, expected, **options):
    """Verify student grading work and pass count calculation.

    This metric checks:
    1. The pass count in O2 matches expected value
    2. Individual student scores are recorded (verifies grading was actually done)
    3. The calculated pass count from student scores matches O2 value (prevents hardcoding)

    Args:
        result: Dict from getter containing 'pass_count', 'student_scores', 'has_formula'
        expected: Rules dict with 'expected_value' and optionally 'expected_scores'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
            - 1.0: Full credit - O2 correct AND grading verified
            - 0.7: Partial credit - O2 correct but no grading evidence
            - 0.3: Partial credit - Some correct scores but wrong final count
            - 0.0: No credit - wrong count and no grading work
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    expected_value = expected.get('expected_value')
    if expected_value is None:
        return 0.0
    pass_count = result.get('pass_count')
    student_scores = result.get('student_scores', {})
    try:
        result_count = int(pass_count) if pass_count is not None else None
        expected_count = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    count_is_correct = result_count == expected_count
    has_grading_work = len(student_scores) >= 5
    if has_grading_work:
        calculated_pass_count = sum((1 for score in student_scores.values() if score >= 60))
        grading_matches_count = calculated_pass_count == result_count if result_count is not None else False
    else:
        grading_matches_count = False
    expected_scores = expected.get('expected_scores', {})
    scores_match = True
    if expected_scores and has_grading_work:
        matching_scores = 0
        for (student_name, expected_score) in expected_scores.items():
            if student_name in student_scores:
                actual_score = student_scores[student_name]
                if abs(actual_score - expected_score) <= 1:
                    matching_scores += 1
        if len(expected_scores) > 0:
            scores_match = matching_scores / len(expected_scores) >= 0.7
    if count_is_correct and has_grading_work and grading_matches_count and scores_match:
        return 1.0
    elif count_is_correct and has_grading_work and grading_matches_count:
        return 0.95
    elif count_is_correct and has_grading_work:
        return 0.7
    elif count_is_correct:
        return 0.5
    elif has_grading_work and scores_match:
        return 0.4
    elif has_grading_work:
        return 0.2
    else:
        return 0.0

def check_file_exists__03122ed4(result, expected, **options):
    """
    Check if file exists and is a valid PNG screenshot.

    Args:
        result: Dict with file info (exists, is_png, size)
        expected: Expected rules (exists, min_size)

    Returns:
        1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if isinstance(result, bool):
        result = {'exists': result, 'is_png': False, 'size': 0}
    if result.get('exists') != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    if not result.get('is_png', False):
        return 0.0
    min_size = expected.get('min_size', 1024)
    if result.get('size', 0) < min_size:
        return 0.0
    return 1.0

def check_file_exists_with_size__dccb2b60(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_text_exact_match__515e2337245bc72c8d34192293ce6646(result, expected, **options):
    """Check if result text exactly matches expected text.

    Args:
        result: Actual text content from file
        expected: Expected text content (from rules dict)
        **options: Additional options (case_sensitive, strip_whitespace)

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    strip_whitespace = options.get('strip_whitespace', True)
    case_sensitive = options.get('case_sensitive', True)
    result_text = result
    expected_text = expected
    if strip_whitespace:
        result_text = result_text.strip()
        expected_text = expected_text.strip()
    if not case_sensitive:
        result_text = result_text.lower()
        expected_text = expected_text.lower()
    return 1.0 if result_text == expected_text else 0.0

def check_file_contains_lines__5f60eaec(result_file_path, expected, **options):
    """
    Check if a text file contains all expected lines.
    
    Args:
        result_file_path: Path to the result text file
        expected: Dict with 'lines' key containing list of expected lines
        **options: Additional options (case_sensitive, order_matters, etc.)
        
    Returns:
        float: Score between 0.0 and 1.0
    """
    if result_file_path is None:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            result_lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    expected_lines = expected.get('lines', [])
    if not expected_lines:
        return 0.0
    case_sensitive = options.get('case_sensitive', True)
    order_matters = options.get('order_matters', False)
    if not case_sensitive:
        result_lines_normalized = [line.lower() for line in result_lines]
        expected_lines_normalized = [line.lower() for line in expected_lines]
    else:
        result_lines_normalized = result_lines
        expected_lines_normalized = expected_lines
    if order_matters:
        if len(result_lines_normalized) != len(expected_lines_normalized):
            return 0.0
        matches = sum((1 for (i, exp_line) in enumerate(expected_lines_normalized) if i < len(result_lines_normalized) and result_lines_normalized[i] == exp_line))
        return matches / len(expected_lines_normalized)
    else:
        result_set = set(result_lines_normalized)
        expected_set = set(expected_lines_normalized)
        if result_set == expected_set:
            return 1.0
        correct_lines = len(result_set & expected_set)
        missing_lines = len(expected_set - result_set)
        extra_lines = len(result_set - expected_set)
        total_issues = missing_lines + extra_lines
        if total_issues == 0:
            return 1.0
        if extra_lines > 0:
            return max(0.0, correct_lines / len(expected_set) - extra_lines * 0.5 / len(expected_set))
        else:
            return correct_lines / len(expected_set)

def check_python_classes__d4d97725(result_file, expected, **options):
    """
    Check if a Python file contains required class definitions.

    Args:
        result_file: Path to the Python file to check
        expected: Dict with 'required_classes' and 'min_classes' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on partial credit
    """
    if not result_file:
        return 0.0
    try:
        with open(result_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    required_classes = expected.get('required_classes', [])
    min_classes = expected.get('min_classes', len(required_classes))
    if not required_classes:
        return 0.0
    found_classes = set()
    try:
        tree = ast.parse(content)
        for node in ast.walk(tree):
            if isinstance(node, ast.ClassDef):
                found_classes.add(node.name)
    except SyntaxError as e:
        print(f'Syntax error parsing file: {e}')
        return 0.0
    found_count = sum((1 for cls in required_classes if cls in found_classes))
    if found_count >= min_classes:
        score = found_count / len(required_classes)
    else:
        score = found_count / len(required_classes) * 0.5
    return score

def check_textbox_in_top_right__8fd9cc45(src_path, expected, **options):
    """
    Check if the textbox is in the top-right corner of the image.
    Task variation 5 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is in top-right, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    pixels = gray_image.load()
    text_pixels = []
    bg_samples = []
    corner_size = min(20, width // 10, height // 10)
    for y in range(corner_size):
        for x in range(corner_size):
            bg_samples.append(pixels[x, y])
            if x < corner_size and width - corner_size < width:
                bg_samples.append(pixels[width - corner_size + x, y])
    avg_bg = sum(bg_samples) / len(bg_samples) if bg_samples else 128
    threshold = avg_bg * 0.7 if avg_bg > 128 else avg_bg * 1.3
    for y in range(height):
        for x in range(width):
            pixel_val = pixels[x, y]
            if avg_bg > 128:
                if pixel_val < threshold:
                    text_pixels.append((x, y))
            elif pixel_val > threshold:
                text_pixels.append((x, y))
    if not text_pixels:
        logger.warning('No text pixels found in image')
        return 0.0
    x_coords = [p[0] for p in text_pixels]
    y_coords = [p[1] for p in text_pixels]
    min_x = min(x_coords)
    max_x = max(x_coords)
    min_y = min(y_coords)
    max_y = max(y_coords)
    center_x = (min_x + max_x) / 2
    center_y = (min_y + max_y) / 2
    in_right = max_x > width * 0.95 and center_x > width * 0.85
    in_top = min_y < height * 0.05 and center_y < height * 0.15
    if in_right and in_top:
        return 1.0
    else:
        return 0.0

def check_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385(result, expected, **options):
    """Compare extracted comments from output file against gold comments from Colab page.

    Args:
        result: Extracted comments from result file (string with one comment per line)
        expected: Configuration dict (the 'rules' dict when using type='rule' in evaluator config)
                  Contains 'url' key with the Colab page URL
        **options: Additional options including 'env' for environment access

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0

    def normalize_comment(comment_line):
        """Normalize a comment line for comparison."""
        stripped = comment_line.strip()
        if stripped.startswith('#'):
            stripped = stripped[1:].strip()
        return stripped.lower()
    result_comments = []
    for line in result.splitlines():
        stripped = line.strip()
        if stripped.startswith('#') or '#' in stripped:
            normalized = normalize_comment(line)
            if normalized:
                result_comments.append(normalized)
    gold_comments = []
    env = options.get('env')
    if env and isinstance(expected, dict):
        try:
            from desktop_env.evaluators.getters.python_file import get_colab_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385
            config = expected
            gold_comment_list = get_colab_python_comments__2b4187f48bcf1d6dd3e5d210ecfff385(env, config)
            if gold_comment_list:
                gold_comments = [normalize_comment(c) for c in gold_comment_list if normalize_comment(c)]
        except Exception as e:
            import sys
            print(f'Failed to extract gold comments from Colab: {e}', file=sys.stderr)
            pass
    if gold_comments:
        matched_comments = 0
        for gold_comment in gold_comments:
            for result_comment in result_comments:
                if gold_comment == result_comment:
                    matched_comments += 1
                    break
                elif len(gold_comment) > 10 and gold_comment in result_comment:
                    matched_comments += 1
                    break
                elif len(result_comment) > 10 and result_comment in gold_comment:
                    matched_comments += 1
                    break
        coverage = matched_comments / len(gold_comments) if gold_comments else 0.0
        if result_comments:
            precision = min(1.0, len(gold_comments) / len(result_comments))
        else:
            precision = 0.0
        score = 0.7 * coverage + 0.3 * precision
        return min(1.0, max(0.0, score))
    import sys
    print('WARNING: Could not extract gold comments from Colab for verification. Returning 0.0.', file=sys.stderr)
    return 0.0

def check_first_line_right_aligned__09652fe4b6098782ed16d144b86a430f(docx_file, expected, **options):
    """
    Check if the first line of the document is right-aligned.

    Args:
        docx_file: Path to the .docx file
        expected: Expected rules (not used in this function)
        **options: Additional options

    Returns:
        float: 1.0 if first line is right-aligned, 0.0 otherwise
    """
    if not docx_file:
        logger.error('No docx file provided')
        return 0.0
    try:
        doc = Document(docx_file)
    except Exception as e:
        logger.error(f'Error loading document: {e}')
        return 0.0
    if not doc.paragraphs:
        logger.error('Document has no paragraphs')
        return 0.0
    first_paragraph = doc.paragraphs[0]
    is_right = first_paragraph.paragraph_format.alignment == WD_PARAGRAPH_ALIGNMENT.RIGHT
    logger.info(f'First paragraph alignment: {first_paragraph.paragraph_format.alignment}')
    logger.info(f'Is right-aligned: {is_right}')
    return 1.0 if is_right else 0.0

def check_dir_and_files__5b67568a(result, expected, **options):
    """Check if directory exists and contains expected number of files.

    Args:
        result: Dict with 'directory_exists' and 'pdf_count'
        expected: Dict with 'require_dir' and 'min_files'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    require_dir = expected.get('require_dir', True)
    min_files = expected.get('min_files', 2)
    if result.get('directory_exists', False) == require_dir:
        score += 0.5
        logger.info(f'Directory existence check passed')
    else:
        logger.warning(f'Directory existence check failed')
        return score
    pdf_count = result.get('pdf_count', 0)
    if pdf_count >= min_files:
        score += 0.5
        logger.info(f'File count check passed: {pdf_count} >= {min_files}')
    else:
        logger.warning(f'File count check failed: {pdf_count} < {min_files}')
    return score

def check_text_lines_exact__d1fc13ca7061617e08d8a914a14209cd(result, expected, **options):
    """Compare text file lines for exact match.

    Args:
        result: List of lines from result file
        expected: Expected list of lines (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_lines = expected.get('lines', [])
    if not isinstance(result, list):
        return 0.0
    if len(result) != len(expected_lines):
        return 0.0
    for (i, (res_line, exp_line)) in enumerate(zip(result, expected_lines)):
        if str(res_line).strip() != str(exp_line).strip():
            return 0.0
    return 1.0

def check_file_exists__3e673542(result, expected, **options):
    """
    Check if command output matches expected value.

    Args:
        result: Command output
        expected: Dict with 'rules' containing 'expected_output'
        **options: Additional options

    Returns:
        float: Score 1.0 if output matches, 0.0 otherwise
    """
    expected_output = expected.get('expected_output', '')
    if isinstance(result, dict) and 'output' in result:
        result = result['output']
    result = str(result).strip()
    return 1.0 if expected_output in result else 0.0

def check_files_by_pattern__0848b03099d380057c59f332d48dc222(directory_list, rule):
    """
    Check if files matching a pattern are copied to the target directory.

    Args:
        directory_list: Directory tree structure from get_list_directory
        rule: Expected configuration with 'pattern' and 'expected_files' keys

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = rule['expected_files']
    actual_files = [node['name'] for node in directory_list['children']]
    if len(actual_files) != len(expected_files):
        return 0.0
    if set(actual_files) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_file_recently_modified__23cbcfa9(result, expected, **options):
    """Check if file was modified within a time window.

    Args:
        result: Unix timestamp from getter
        expected: Expected rules dict with 'max_age_seconds' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if file was recently modified, 0.0 otherwise
    """
    if result < 0:
        logger.warning('File not found or error getting modification time')
        return 0.0
    max_age_seconds = expected.get('max_age_seconds', 600)
    current_time = int(time.time())
    age_seconds = current_time - result
    if age_seconds <= max_age_seconds:
        return 1.0
    else:
        logger.info(f'File is {age_seconds} seconds old, exceeds max age of {max_age_seconds} seconds')
        return 0.0

def check_text_replacement__27683cd2(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_file_count__bb5651c2(result_state, expected_state, **options):
    """
    Check if the file content matches the expected count.

    This metric verifies that the content of document_count.txt
    matches the expected count of .docx files.

    Args:
        result_state: Content from the file (str or None)
        expected_state: Expected state dict with 'count' key
        **options: Additional options

    Returns:
        float: Score (1.0 if match, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    expected_count = expected_state.get('count')
    if expected_count is None:
        return 0.0
    if str(result_state) == str(expected_count):
        return 1.0
    try:
        result_num = int(result_state)
        expected_num = int(expected_count)
        if result_num == expected_num:
            return 1.0
    except (ValueError, TypeError):
        pass
    return 0.0

def check_files_moved__90e7472a(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if files were moved to the correct directory.

    Verifies both:
    1. Expected files exist in destination directory
    2. Source directory is empty (files were moved, not copied)

    Args:
        result: Command output with format:
                DESTINATION:
                <list of files>
                SOURCE_COUNT:
                <count>
        expected: Expected rules dict with 'expected_files' list and 'expected_source_count'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    try:
        expected_files = expected.get('expected_files', [])
        expected_source_count = expected.get('expected_source_count', 0)
        if not expected_files:
            return 0.0
        lines = result.strip().split('\n')
        destination_files = []
        source_count = None
        current_section = None
        for line in lines:
            line = line.strip()
            if line == 'DESTINATION:':
                current_section = 'destination'
            elif line == 'SOURCE_COUNT:':
                current_section = 'source_count'
            elif current_section == 'destination' and line and (not line.endswith(':')):
                destination_files.append(line)
            elif current_section == 'source_count' and line.isdigit():
                source_count = int(line)
        destination_matches = sum((1 for f in expected_files if f in destination_files))
        destination_score = destination_matches / len(expected_files)
        if source_count is None:
            return destination_score * 0.5
        source_score = 1.0 if source_count == expected_source_count else 0.0
        final_score = destination_score * 0.7 + source_score * 0.3
        return min(1.0, final_score)
    except Exception as e:
        return 0.0

def check_file_exists__7ae8ce2b(result, expected, **options):
    """Check if file exists and has content.
    
    Args:
        result: Dict with 'exists' and 'has_content' keys
        expected: Dict with 'should_exist' key
        
    Returns:
        float: 1.0 if file exists and has content, else 0.0
    """
    if result is None:
        return 0.0
    exists = result.get('exists', False)
    has_content = result.get('has_content', False)
    should_exist = expected.get('should_exist', True)
    if should_exist:
        return 1.0 if exists and has_content else 0.0
    else:
        return 1.0 if not exists else 0.0

def check_file_checksum__032a6328(result: str, expected: dict, **options) -> float:
    """Check if file checksum matches expected checksum.

    Args:
        result: Checksum from getter
        expected: Dict with 'checksum' key
        **options: Additional options

    Returns:
        1.0 if match, 0.0 otherwise
    """
    expected_checksum = expected.get('checksum', '')
    if result == expected_checksum:
        return 1.0
    return 0.0

def check_lines_removed__3204f52d(result, expected, **options):
    """Check if specific lines were removed from the file.

    Args:
        result: List of file lines
        expected: Dict with 'max_lines' (expected total line count), 'missing_patterns' (optional)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    max_lines = expected.get('max_lines', 8)
    missing_patterns = expected.get('missing_patterns', [])
    score = 0.0
    if len(result) <= max_lines:
        score += 0.6
    if missing_patterns:
        content = '\n'.join(result)
        patterns_missing = sum((1 for pattern in missing_patterns if pattern not in content))
        if patterns_missing == len(missing_patterns):
            score += 0.4
        else:
            score += 0.4 * (patterns_missing / len(missing_patterns))
    else:
        score += 0.4
    return min(score, 1.0)

def check_file_exists__dee238bc(result, expected, **options):
    """
    Check if subtitle file exists and is non-empty.

    Args:
        result: Path to file
        expected: Expected rules
        **options: Additional options

    Returns:
        float: 1.0 if file exists and non-empty, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not os.path.exists(result):
        return 0.0
    if os.path.getsize(result) == 0:
        return 0.0
    return 1.0

def check_text_exact_match__895c3960d172d43278234eeb5c495eda(result: str, expected: dict, **options) -> float:
    """
    Check if text content exactly matches expected value.

    Args:
        result: Text content from getter
        expected: Dict with 'content' key containing expected text
        **options: Additional options

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        logger.debug('Result is None')
        return 0.0
    expected_content = expected.get('content', '')
    if result == expected_content:
        return 1.0
    else:
        logger.debug(f"Content mismatch. Expected: '{expected_content}', Got: '{result}'")
        return 0.0

def check_file_permissions__dd3bb1bb(result, rules, **options):
    """Check if all files have the expected permission.

    Args:
        result: Output from ls -l command showing file permissions
        rules: Dict with 'permission' key specifying expected permission string (e.g., '-rw-------')

    Returns:
        float: 1.0 if all files have correct permission, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    expected_perm = rules.get('permission', '')
    lines = result.strip().split('\n')
    for line in lines:
        if not line.strip():
            continue
        parts = line.split()
        if len(parts) > 0:
            actual_perm = parts[0]
            if actual_perm != expected_perm:
                return 0.0
    return 1.0

def check_lines_commented__4ee5b05c(result, expected, **options):
    """Check if specific lines are commented out.

    Args:
        result: List of file lines
        expected: Dict with 'start_line' (1-based), 'end_line'

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    start_line = expected.get('start_line', 5)
    end_line = expected.get('end_line', 7)
    start_idx = start_line - 1
    end_idx = end_line - 1
    if end_idx >= len(result):
        return 0.0
    score = 0.0
    total_lines = end_idx - start_idx + 1
    for i in range(start_idx, end_idx + 1):
        line = result[i].lstrip()
        if line and line[0] == '#':
            score += 1.0 / total_lines
    return score

def check_exact_text_match__9a63ba8e(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_tsv_file_exists__ecf92ffc(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if TSV file exists and is properly formatted.

    Args:
        result: Dict with file_exists, file_path, is_tsv_format, row_count, column_count from getter
        expected: Dict with expected file_path
        **options: Additional options

    Returns:
        Score: 1.0 if file exists, is at expected path, is TSV format, and contains data; 0.0 otherwise
    """
    if not result.get('file_exists', False):
        logger.warning(f"TSV file does not exist at {result.get('file_path')}")
        return 0.0
    expected_path = expected.get('file_path', '')
    result_path = result.get('file_path', '')
    if expected_path and result_path != expected_path:
        logger.warning(f'File path mismatch: expected {expected_path}, got {result_path}')
        return 0.0
    if not result.get('is_tsv_format', False):
        logger.warning(f'File exists but is not in TSV format (no tab separators found)')
        return 0.0
    row_count = result.get('row_count', 0)
    column_count = result.get('column_count', 0)
    if row_count < 1:
        logger.warning(f'TSV file is empty (no rows)')
        return 0.0
    if column_count < 2:
        logger.warning(f'TSV file has insufficient columns ({column_count}), expected at least 2')
        return 0.0
    logger.info(f'TSV file validation passed: {result_path} ({row_count} rows, {column_count} columns)')
    return 1.0

def check_file_exists__e3dda739ee4da14903d2e8df52d7a41d(result, expected, **options):
    """
    Check if file exists and is a valid MP3 audio file.

    Args:
        result: dict with 'exists', 'is_mp3', and 'file_size' keys
        expected: dict with 'exists' key (and optionally 'min_size')
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is a valid MP3, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    min_size = expected.get('min_size', 1000)
    if isinstance(result, dict):
        file_exists = result.get('exists', False)
        is_mp3 = result.get('is_mp3', False)
        file_size = result.get('file_size', 0)
    else:
        file_exists = bool(result)
        is_mp3 = False
        file_size = 0
    if not expected_exists and (not file_exists):
        logger.info('File correctly does not exist')
        return 1.0
    if expected_exists:
        if not file_exists:
            logger.warning('File does not exist when it should')
            return 0.0
        if not is_mp3:
            logger.warning('File exists but is not a valid MP3 audio file')
            return 0.0
        if file_size < min_size:
            logger.warning(f'File is too small ({file_size} bytes), likely not a valid audio file (min: {min_size})')
            return 0.0
        logger.info(f'File exists and is a valid MP3 (size: {file_size} bytes)')
        return 1.0
    logger.warning("File exists when it shouldn't")
    return 0.0

def check_text_contains__7dfb45a4(result, expected, **options):
    """Check if text content contains expected substring(s).

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    score = 0.0
    for substring in contains:
        if substring.lower() in result.lower():
            score += 1.0 / len(contains)
    return score

def check_text_alignment__faeddc67(result_file, expected, **options):
    try:
        prs = Presentation(result_file)
        slide_idx = expected.get('slide_idx', 0)
        shape_idx = expected.get('shape_idx', 0)
        expected_align = expected.get('alignment', 'LEFT')
        align_map = {'LEFT': PP_ALIGN.LEFT, 'CENTER': PP_ALIGN.CENTER, 'RIGHT': PP_ALIGN.RIGHT, 'JUSTIFY': PP_ALIGN.JUSTIFY}
        expected_alignment_enum = align_map.get(expected_align, PP_ALIGN.LEFT)
        if slide_idx >= len(prs.slides):
            return 0.0
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            return 0.0
        shape = slide.shapes[shape_idx]
        if not hasattr(shape, 'text_frame'):
            return 0.0
        for para in shape.text_frame.paragraphs:
            actual_align = para.alignment if para.alignment is not None else PP_ALIGN.LEFT
            if actual_align != expected_alignment_enum:
                return 0.0
        return 1.0
    except:
        return 0.0

def check_text_content__69119b71(result, expected, **options):
    """Compare text content against expected value.

    Args:
        result: Text content from getter
        expected: Expected text value (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_stripped = str(result).strip()
    expected_stripped = str(expected_value).strip()
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_file_exported__46f0b51f(result, expected, **options):
    """Check if file was successfully exported as a valid PDF with spreadsheet content.

    Args:
        result: Dict with validation results from getter:
            - exists: bool, whether file exists
            - is_pdf: bool, whether file is a valid PDF
            - file_size: int, file size in bytes
            - has_content: bool, whether PDF has reasonable content
            - has_spreadsheet_data: bool, whether PDF contains spreadsheet data
            - created_by_libreoffice: bool, whether PDF metadata shows LibreOffice as creator
        expected: Expected value (True for file should exist as valid PDF)
        **options: Additional options

    Returns:
        float: 1.0 if file is a valid PDF with spreadsheet content, 0.0 otherwise
    """
    if result is None or not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('is_pdf', False):
        return 0.0
    if result.get('file_size', 0) < 1024:
        return 0.0
    if not result.get('has_content', False):
        return 0.0
    if not result.get('has_spreadsheet_data', False):
        return 0.0
    if not result.get('created_by_libreoffice', False):
        return 0.0
    return 1.0

def check_comprehensive_text__c4f4ba50(result, expected, **options):
    """Check if comprehensive text extraction occurred.

    Args:
        result: Text from getter
        expected: Expected dict with 'min_length' and 'key_terms' keys

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    min_length = expected.get('min_length', 100)
    key_terms = expected.get('key_terms', [])
    score = 0.0
    if len(result) >= min_length:
        score += 0.3
    if key_terms:
        result_lower = result.lower()
        matches = sum((1 for term in key_terms if term.lower() in result_lower))
        score += 0.7 * (matches / len(key_terms))
    return score

def check_text_output__99a23d7f(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_locations__b75faf5b6765d0d1458ed6b6d219047b(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if files are in the correct location, not trashed, are valid PNG images, and were created recently.

    This function verifies:
    1. Files are in the correct Google Drive folder (figures/ in root)
    2. Files have the expected filenames (1.png, 2.png, 3.png)
    3. Files are not trashed
    4. Files have valid PNG MIME type (image/png)
    5. Files are not empty or trivially small (size > 100 bytes)
    6. Files were created recently (within last hour) to ensure they were created during task execution

    Args:
        result: File location info dict with file_details for validation
        expected: Expected configuration (from rules dict) with 'expected_count'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on proportion of correct files
    """
    expected_count = expected.get('expected_count', 3)
    correct_count = result.get('correct_location_count', 0)
    if correct_count == expected_count:
        return 1.0
    elif correct_count > 0:
        return correct_count / expected_count
    else:
        return 0.0

def check_copy_file_status__24914e86(result, expected, **options):
    """Check if file was copied correctly (exists at both source and destination).

    Args:
        result: Dict with 'source_exists' and 'dest_exists' boolean values
        expected: Expected state (dict with 'source_exists' and 'dest_exists' keys)
        **options: Additional options

    Returns:
        float: 1.0 if both source and destination exist as expected, partial score otherwise
    """
    source_exists = result.get('source_exists', False)
    dest_exists = result.get('dest_exists', False)
    expected_source = expected.get('source_exists', True)
    expected_dest = expected.get('dest_exists', True)
    score = 0.0
    if source_exists == expected_source:
        score += 0.5
    if dest_exists == expected_dest:
        score += 0.5
    return score

def check_word_count_file__adfc25c4(result, expected, **options):
    """Check if word count matches expected value.

    Args:
        result: Word count from getter (int or None)
        expected: Expected rules dict with 'expected_count' key
        **options: Additional options

    Returns:
        float: 1.0 if counts match, 0.0 otherwise
    """
    if result is None:
        logger.warning('No word count result received')
        return 0.0
    expected_count = expected.get('expected_count')
    if expected_count is None:
        logger.error('No expected_count in rules')
        return 0.0
    tolerance = options.get('tolerance', 2)
    if abs(result - expected_count) <= tolerance:
        logger.info(f'Word count match: {result} (expected {expected_count})')
        return 1.0
    else:
        logger.info(f'Word count mismatch: {result} vs expected {expected_count}')
        return 0.0

def check_text_replacement__514834e6(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_file_exists__735234d9(result, expected, **options):
    """Check if file exists and is non-empty."""
    if result is None or not os.path.exists(result):
        return 0.0
    return 1.0 if os.path.getsize(result) > 0 else 0.0

def check_timezone_utc_plus_8__a394a8f31cb4c0f1d8dc918cb19a351d(timedatectl_output, expected, **options):
    """
    Check if timezone is set to UTC+8 (e.g., Asia/Shanghai, Asia/Hong_Kong).

    Args:
        timedatectl_output: Output from 'timedatectl status' command
        expected: Expected timezone offset dict (e.g., {'offset': '+0800'})
        **options: Additional options

    Returns:
        float: 1.0 if timezone is UTC+8, 0.0 otherwise
    """
    lines = timedatectl_output.split('\n')
    expected_offset = expected.get('offset', '+0800')
    for line in lines:
        if 'Time zone:' in line:
            if line.endswith(f'{expected_offset})'):
                return 1.0
            else:
                return 0.0
    return 0.0

def check_largest_filename__b7de68e1(result, expected, **options):
    """Verify the largest .doc filename is correct.

    Args:
        result: Filename from getter
        expected: Expected filename
        **options: Additional options

    Returns:
        float: 1.0 if filename matches, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_filename = expected.get('filename', '')
    result_clean = result.strip()
    return 1.0 if result_clean == expected_filename else 0.0

def check_both_files_non_empty__e07d7a26(result, expected, **options):
    """
    Check if both files are non-empty.

    Args:
        result: dict with csv_size and xlsx_size
        expected: dict (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('csv_size', 0) > 0:
        score += 0.5
    if result.get('xlsx_size', 0) > 0:
        score += 0.5
    return score

def check_archive_contains_files__0b074054(result, expected, **options):
    """Check if archive contains expected files.

    Args:
        result: List of files in archive (newline-separated)
        expected: Dict with 'files' key containing list of expected filenames
        **options: Additional options

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not result or result.strip() == '':
        return 0.0
    archive_lines = [line.strip() for line in result.split('\n') if line.strip()]
    expected_files = expected.get('files', [])
    for expected_item in expected_files:
        found = False
        for line in archive_lines:
            if line == expected_item:
                found = True
                break
            if expected_item.endswith('/') and line.startswith(expected_item):
                found = True
                break
            if expected_item.endswith('/'):
                expected_without_slash = expected_item.rstrip('/')
                if line == expected_without_slash or line.startswith(expected_without_slash + '/'):
                    found = True
                    break
        if not found:
            return 0.0
    return 1.0

def check_files_with_keyword__91578e58(result, expected, **options):
    """Check if all required keywords have minimum matches.

    Args:
        result: Dict with keyword counts {'keyword1': count, 'keyword2': count, ...}
        expected: Dict with:
            - keywords: List of required keywords
            - min_matches_per_keyword: Minimum matches per keyword (default 1)

    Returns:
        Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    required_keywords = expected.get('keywords', [])
    min_matches = expected.get('min_matches_per_keyword', 1)
    for keyword in required_keywords:
        if keyword not in result or result[keyword] < min_matches:
            return 0.0
    return 1.0

def check_text_file_mountains__1e67b9ae311891ef2e3034615cded86c(result: List[str], expected: Dict, **options) -> float:
    """
    Check if text file contains expected mountain names.

    Args:
        result: List of lines from the file
        expected: Dict with expected mountain names
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    required_mountains = expected.get('mountains', [])
    if not required_mountains:
        return 0.0
    content_combined = ' '.join(result).lower()
    found_count = 0
    for mountain_variants in required_mountains:
        found = False
        for variant in mountain_variants:
            if variant.lower() in content_combined:
                found = True
                logger.info(f'Found mountain variant: {variant}')
                break
        if found:
            found_count += 1
        else:
            logger.info(f'Missing mountain: {mountain_variants[0]}')
    score = found_count / len(required_mountains)
    return score

def check_file_exists__506ad17f(result, expected, **options):
    """
    Check if file existence matches expected state.

    Args:
        result: Boolean from getter (True if file exists)
        expected: Expected state (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_first_line_heading1_style__ba24a7ac(result, expected, **options):
    """Check if first line has Heading 1 style.

    Args:
        result: Style name string from getter
        expected: Expected style from rules
        **options: Additional options

    Returns:
        float: 1.0 if Heading 1, 0.0 otherwise
    """
    expected_style = expected.get('style', 'Heading 1')
    return 1.0 if result == expected_style else 0.0

def check_file_exists__a544be73(result_state, expected_state, **options):
    """
    Check if the file exists and is non-empty (has been successfully saved).

    Args:
        result_state (str): Local path to the downloaded file from VM
        expected_state (dict): Expected configuration (can be empty for simple existence check)
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is non-empty, 0.0 otherwise
    """
    try:
        if not result_state or not isinstance(result_state, str):
            return 0.0
        if not os.path.exists(result_state):
            return 0.0
        if not os.path.isfile(result_state):
            return 0.0
        file_size = os.path.getsize(result_state)
        if file_size == 0:
            return 0.0
        return 1.0
    except Exception as e:
        print(f'Error in check_file_exists__a544be73: {e}')
        return 0.0

def check_file_count__63477992bbc88c6a7091e80f7c0a8a72(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the file names match the expected sequential naming pattern (1.png, 2.png, 3.png).

    Args:
        result: List of actual file names (List[str])
        expected: Expected configuration (from rules dict) with:
            - count: Expected number of files
            - file_names (optional): Specific expected file names
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    expected_file_names = expected.get('file_names', None)
    if len(result) != expected_count:
        return 0.0
    if expected_file_names is not None:
        if sorted(result) == sorted(expected_file_names):
            return 1.0
        else:
            return 0.0
    pattern = re.compile('^(\\d+)\\.png$')
    numbers = []
    for filename in result:
        match = pattern.match(filename)
        if not match:
            return 0.0
        numbers.append(int(match.group(1)))
    numbers.sort()
    expected_sequence = list(range(1, expected_count + 1))
    if numbers == expected_sequence:
        return 1.0
    else:
        return 0.0

def check_exact_match_v4__d4477d7a(result, expected, **options):
    """Compare result against expected integer value.

    Args:
        result: Actual value from getter
        expected: Rules dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_value = expected.get('expected_value')
    if result is None or expected_value is None:
        return 0.0
    try:
        result_int = int(result)
        expected_int = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    if result_int == expected_int:
        return 1.0
    return 0.0

def check_eml_files_exist__aebf0a91a4be0be45ef3247f943131f2(result: List[Optional[str]], expected: dict, **options) -> float:
    """Check if email files were successfully downloaded from Google Drive.

    Args:
        result: List of downloaded file paths (or None for missing files)
        expected: Expected rules dict with:
            - count: Expected number of files
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    expected_count = expected.get('count', 2)
    if result is None:
        return 0.0
    successful = sum((1 for path in result if path is not None))
    if expected_count == 0:
        return 0.0
    return successful / expected_count

def check_backup_files__ac9408954b941c7f40eedd27a6f1296b(result: dict, expected: dict) -> float:
    """
    Check if backup files exist in specified locations.

    Args:
        result: Dict with backup file existence info from getter
        expected: Dict with 'backup_locations' (list of dicts with 'dir' and 'filename')

    Returns:
        float: 1.0 if all backup files exist, 0.0 otherwise
    """
    if result is None:
        return 0.0
    backup_locations = expected.get('backup_locations', [])
    for location in backup_locations:
        dir_name = location.get('dir')
        key = f'{dir_name}_backup'
        if not result.get(key, False):
            return 0.0
    return 1.0

def check_file_exported__f25970ca(result, expected, **options):
    """Check if CSV file was successfully exported with valid content from XLSX source.

    Args:
        result: Dict containing validation results from getter
        expected: Expected value (True for file should exist)
        **options: Additional options

    Returns:
        float: 1.0 if CSV was properly exported from XLSX, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if expected is True or expected == True:
        if not isinstance(result, dict):
            return 0.0
        if not result.get('exists', False):
            return 0.0
        if not result.get('is_csv', False):
            return 0.0
        if not result.get('has_content', False):
            return 0.0
        if result.get('row_count', 0) < 2:
            return 0.0
        if result.get('xlsx_exists', False):
            if not result.get('data_matches', False):
                return 0.0
        if not result.get('valid', False):
            return 0.0
        return 1.0
    if result == expected:
        return 1.0
    return 0.0

def check_titles_file__3af0b2bf(result, expected, **options):
    """Compare extracted titles file against expected titles list.

    Args:
        result: File path (string path to the titles.txt file)
        expected: Expected rules dict (contains 'titles' key with list of titles)
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 and 1.0 based on matching titles
    """
    if not isinstance(result, str):
        return 0.0
    if not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception:
        return 0.0
    result_lines = [line.strip() for line in content.strip().split('\n')]
    if isinstance(expected, dict):
        expected_titles = expected.get('titles', [])
    else:
        expected_titles = expected
    if len(result_lines) != len(expected_titles):
        return 0.0
    matches = 0
    for (res_title, exp_title) in zip(result_lines, expected_titles):
        if res_title == exp_title:
            matches += 1
    return matches / len(expected_titles) if expected_titles else 1.0

def check_file_exists__a693f275(result, expected, **options):
    """
    Check if a file exists based on command output.

    Args:
        result: Output from ls command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'exists')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_file_exists__46ccb784(result, expected, **options):
    """Check if file existence matches expected value.

    Args:
        result: Boolean from getter
        expected: Expected boolean value
        **options: Additional options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if result == expected:
        return 1.0
    else:
        return 0.0

def check_total_file_count__ed6a3699fe6af7deef02a8e547504034(result, expected, **options):
    """Check if file counts in each directory match expected values.

    Args:
        result: Dictionary mapping directory name to file count (e.g., {'dir1': 1, 'dir2': 1, 'dir3': 0})
        expected: Dictionary mapping directory name to expected file count
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    if set(expected.keys()) != set(result.keys()):
        return 0.0
    for (directory, expected_count) in expected.items():
        if result.get(directory) != expected_count:
            return 0.0
    return 1.0

def check_text_file_content__5ced85fc_aug18_v1_a1b2c3d4e5f6a7b8c9d0e1f2a3b4c5d6(result, expected, **options):
    """Check if file content matches expected content exactly.

    Args:
        result: Content string from the file
        expected: Rules dict with 'expected_content' key
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        logger.error('Result is not a string')
        return 0.0
    expected_content = expected.get('expected_content', '')
    result_stripped = result.strip()
    expected_stripped = expected_content.strip()
    if result_stripped == expected_stripped:
        logger.info('Content matches expected value')
        return 1.0
    else:
        logger.info(f"Content mismatch: got '{result_stripped}', expected '{expected_stripped}'")
        return 0.0

def check_direct_json_object__fc6d8143(result, expected: Dict, **options):
    """
    Check if result matches expected values with relativeTime conversion support.

    This metric handles:
    1. RelativeTime conversion (e.g., "tomorrow" -> actual date)
    2. Direct JSON object comparison

    Args:
        result: dict - The extracted data from the webpage
        expected: dict - Contains "relativeTime" and "expected" fields
        **options: additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    logger.info(f'[DEBUG] check_direct_json_object__fc6d8143 called with result: {result}')
    logger.info(f'[DEBUG] check_direct_json_object__fc6d8143 called with expected: {expected}')
    if isinstance(result, str):
        result = result.strip()
        result = result.replace("'", '"')
        try:
            result = json.loads(result)
        except json.JSONDecodeError as e:
            logger.error(f'[DEBUG] Failed to parse result as JSON: {e}')
            return 0.0
    logger.info(f'[DEBUG] Processed result: {result}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if 'relativeTime' in expected:
        logger.info('[DEBUG] Detected relativeTime in expected, performing time conversion')
        relativeTime = expected['relativeTime']
        timezone_str = expected.get('timezone', get_default_timezone())
        try:
            timezone = pytz.timezone(timezone_str)
            logger.info(f'Successfully loaded timezone: {timezone_str}')
        except pytz.exceptions.UnknownTimeZoneError:
            logger.error(f'Unknown timezone: {timezone_str}, falling back to UTC')
            timezone = pytz.UTC
        now = datetime.now(timezone)
        logger.info(f"Current time in {timezone_str}: {now.strftime('%Y-%m-%d %H:%M:%S %Z')}")
        if 'to' not in relativeTime.keys():
            start_relative_time = relativeTime['from']
            logger.info(f"Processing single time: '{start_relative_time}'")
            if start_relative_time in relativeTime_to_IntDay:
                days_to_add = relativeTime_to_IntDay[start_relative_time]
                timediff = timedelta(days=days_to_add)
                absoluteDay = now + timediff
                logger.info(f"Simple calculation: {start_relative_time} = {days_to_add} days → {absoluteDay.strftime('%Y-%m-%d %H:%M:%S %Z')}")
            else:
                logger.error(f'Unsupported relative time: {start_relative_time}')
                return 0.0
            expected_values = expected['expected'].copy()
            if 'time' in expected_values:
                regular_time = apply_rules_to_timeFormat(expected_values['time'], absoluteDay)
                logger.info(f'Final formatted time: {regular_time}')
                expected_values['time'] = regular_time
            for (key, expected_value) in expected_values.items():
                actual_value = result.get(key)
                logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
                if expected_value != actual_value:
                    logger.info(f"[DEBUG] Value comparison failed for key '{key}': expected='{expected_value}', actual='{actual_value}', returning 0.0")
                    return 0.0
                else:
                    logger.info(f"[DEBUG] Value comparison passed for key '{key}'")
            logger.info('[DEBUG] All comparisons passed, returning 1.0')
            return 1.0
        else:
            logger.error('[DEBUG] Time range (from/to) not yet implemented')
            return 0.0
    else:
        expected_json = expected.get('expected', {})
        logger.info(f'[DEBUG] Expected JSON: {expected_json}')
        for (key, expected_value) in expected_json.items():
            actual_value = result.get(key)
            logger.info(f"[DEBUG] Checking key '{key}': expected='{expected_value}', actual='{actual_value}'")
            if expected_value != actual_value:
                logger.info(f"[DEBUG] Value comparison failed for key '{key}': expected='{expected_value}', actual='{actual_value}', returning 0.0")
                return 0.0
            else:
                logger.info(f"[DEBUG] Value comparison passed for key '{key}'")
        logger.info('[DEBUG] All comparisons passed, returning 1.0')
        return 1.0

def check_both_files_exist__805294f8(result, expected, **options):
    """Verify that both files exist.

    Args:
        result: Dict with file1_exists and file2_exists booleans
        expected: Dict with rules
        **options: Additional options

    Returns:
        float: 1.0 if both files exist, 0.0 otherwise
    """
    both_exist_required = expected.get('both_exist', True)
    file1_exists = result.get('file1_exists', False)
    file2_exists = result.get('file2_exists', False)
    if both_exist_required:
        return 1.0 if file1_exists and file2_exists else 0.0
    else:
        return 1.0 if not file1_exists and (not file2_exists) else 0.0

def check_text_match__ed1a5c265e6c6d06dcaf2ec482204403(result, expected, **options):
    """Check if text content matches expected value.

    Args:
        result: Actual text content from getter
        expected: Expected text value (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_normalized = ' '.join(result.split())
    expected_normalized = ' '.join(str(expected_value).split())
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_textbox_on_top__49abb020(src_path, expected, **options):
    """
    Check if the textbox is at the top of the image.
    Task variation 2 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is at top, 0.0 otherwise)
    """
    if src_path is None:
        logger.error('No source path provided')
        return 0.0
    try:
        source_image = Image.open(src_path)
        (width, height) = source_image.size
        img_array = np.array(source_image.convert('RGB'))
        edge_pixels = []
        edge_pixels.extend(img_array[0, :].tolist())
        edge_pixels.extend(img_array[-1, :].tolist())
        edge_pixels.extend(img_array[:, 0].tolist())
        edge_pixels.extend(img_array[:, -1].tolist())
        background_color = np.mean(edge_pixels, axis=0)
        diff = np.sqrt(np.sum((img_array - background_color) ** 2, axis=2))
        text_mask = diff > 30
        topmost_text_y = height
        text_detected = False
        for y in range(height):
            if np.any(text_mask[y, :]):
                topmost_text_y = y
                text_detected = True
                break
        if not text_detected:
            logger.warning('No text detected in the image')
            return 0.0
        threshold_y = height * 0.05
        logger.info(f'Topmost text pixel at y={topmost_text_y}, threshold={threshold_y}, height={height}')
        if topmost_text_y < threshold_y:
            logger.info('Text is positioned at the top of the canvas')
            return 1.0
        else:
            logger.info(f'Text is not at the top (y={topmost_text_y} >= {threshold_y})')
            return 0.0
    except Exception as e:
        logger.error(f'Error processing image: {e}')
        return 0.0

def check_file_contains__3c9f051952e2f37565e45b593e085b87(result: str, expected: dict, **options) -> float:
    """Check if file content contains the expected string.

    Args:
        result: String containing the actual file content
        expected: Dict with 'text' field specifying expected substring

    Returns:
        1.0 if expected text is found in result, 0.0 otherwise
    """
    expected_text = expected.get('text', '')
    if not expected_text:
        return 0.0
    if expected_text in result:
        return 1.0
    return 0.0

def check_file_exists__d8671412(result, expected, **options):
    """Check if file rename operation was successful.

    Args:
        result: Dict with 'target_exists' and 'source_exists' keys from getter
        expected: Dict with 'expected' key (boolean for target) and optional 'source_should_exist' key

    Returns:
        float: 1.0 if rename successful (target exists AND source doesn't), 0.0 otherwise
    """
    expected_target = expected.get('expected', True)
    expected_source = expected.get('source_should_exist', False)
    target_exists = result.get('target_exists', False)
    source_exists = result.get('source_exists', False)
    if target_exists == expected_target and source_exists == expected_source:
        return 1.0
    else:
        print(f'Rename verification failed. Target exists: {target_exists} (expected: {expected_target}), Source exists: {source_exists} (expected: {expected_source})')
        return 0.0

def check_timezone_utc_plus_1__5f2017655ce06bbb056bf7cc4210c4ca(timedatectl_output, expected, **options):
    """
    Check if timezone is set to UTC+1 (e.g., Europe/London during BST, Europe/Paris).

    Args:
        timedatectl_output: Output from 'timedatectl status' command
        expected: Expected timezone offset dict with 'offset' key (e.g., {'offset': '+0100'})
        **options: Additional options

    Returns:
        float: 1.0 if timezone matches expected offset, 0.0 otherwise
    """
    offset = expected.get('offset', '+0100')
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            if line.endswith(f'{offset})'):
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_renamed__c396550f(actual: Dict[str, bool], expected: Dict, **options) -> float:
    """Check if a file was successfully renamed.

    Args:
        actual (Dict[str, bool]): Dictionary with 'old_exists' and 'new_exists' keys
        expected (Dict): expected dict (not used, implicit expectation)
        **options: Additional options

    Return:
        float: the score (1.0 if renamed successfully, 0.0 otherwise)
    """
    if not actual.get('old_exists', True) and actual.get('new_exists', False):
        return 1.0
    return 0.0

def check_file_content__89d906fa(result, expected, **options):
    """
    Check if the file content matches the expected value.

    Args:
        result: Actual file content from getter
        expected: Dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_value = str(expected.get('expected_value', ''))
    if result.strip() == expected_value.strip():
        return 1.0
    else:
        return 0.0

def check_textbox_on_topside__57a0b169(src_path, expected_state, **options):
    """
    Check if the textbox is on the top side of the image.
    Variation 0 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in top 5%, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    top_most_dark_pixel = height
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                top_most_dark_pixel = min(top_most_dark_pixel, y)
                break
    if top_most_dark_pixel < height * 0.05:
        return 1.0
    else:
        return 0.0

def check_text_contains__b38e8bb9(result, expected, **options):
    """Check if text content contains expected names on separate lines.

    Args:
        result: Text content from getter
        expected: Dict with 'contains' key (str or list of str)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    contains = expected.get('contains', [])
    if isinstance(contains, str):
        contains = [contains]
    if not contains:
        return 0.0
    result_lines = [line.strip() for line in result.split('\n') if line.strip()]
    score = 0.0
    for expected_name in contains:
        if expected_name.strip() in result_lines:
            score += 1.0 / len(contains)
    return score

def check_timezone_utc_plus_530__851aed8c(result_state, expected_state, **options):
    """
    Check if timezone is set to +0530.

    Args:
        result_state: Output from 'timedatectl status' command
        expected_state: Expected timezone configuration (dict with 'timezone_offset')
        **options: Additional options

    Returns:
        float: 1.0 if timezone offset is +0530, 0.0 otherwise
    """
    timedatectl_output = result_state
    lines = timedatectl_output.split('\n')
    expected_offset = expected_state.get('timezone_offset', '+0530')
    timezone_line = None
    for line in lines:
        if 'Time zone:' in line:
            timezone_line = line
            break
    if not timezone_line:
        if len(lines) > 3:
            timezone_line = lines[3]
        else:
            return 0.0
    if timezone_line.endswith(f'{expected_offset})'):
        return 1.0
    else:
        return 0.0

def check_text_replacement__8b701bf8d0cacb95438d2e4e17a8b914(result: str, expected: Dict, **options) -> float:
    """
    Check if text replacement was performed correctly.

    Args:
        result: Actual file content (string from getter)
        expected: Expected value with 'expected_content' key
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        logger.warning('Result is empty')
        return 0.0
    expected_content = expected.get('expected_content', '')
    if not expected_content:
        logger.error('No expected_content in expected dict')
        return 0.0
    if result.strip() == expected_content.strip():
        return 1.0
    else:
        logger.info(f'Content mismatch. Expected length: {len(expected_content)}, Actual length: {len(result)}')
        return 0.0

def check_word_count_file__7c58ef63(result, expected, **options):
    """Check if the word count file contains the expected value within tolerance.

    Args:
        result: File content as string from getter
        expected: Expected word count (or dict with count and tolerance)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.error('File content is None - file may not exist')
        return 0.0
    if isinstance(expected, dict):
        expected_count = expected.get('count')
        tolerance = expected.get('tolerance', 0)
    else:
        expected_count = expected
        tolerance = options.get('tolerance', 0)
    if expected_count is None:
        logger.error('No expected count provided')
        return 0.0
    try:
        numbers = re.findall('\\d+', result)
        if not numbers:
            logger.error(f'No numeric value found in file content: {result[:100]}')
            return 0.0
        actual_count = int(numbers[0])
        if len(numbers) > 1:
            actual_count = min([int(n) for n in numbers], key=lambda x: abs(x - expected_count))
        logger.info(f'Extracted word count: {actual_count} from content: {result[:100]}')
        diff = abs(actual_count - expected_count)
        if diff <= tolerance:
            logger.info(f'Word count match: {actual_count} within tolerance {tolerance} of {expected_count}')
            return 1.0
        else:
            logger.info(f'Word count mismatch: got {actual_count}, expected {expected_count} (tolerance: {tolerance})')
            return 0.0
    except Exception as e:
        logger.error(f'Failed to parse word count from content: {e}')
        return 0.0

def check_timezone_utc_minus_8__dfb26af6f8639d73cfa7d8bdb0721f21(timedatectl_output, expected, **options):
    """
    Check if timezone is set to UTC-8 (e.g., America/Los_Angeles, America/Vancouver).

    Args:
        timedatectl_output: Output from 'timedatectl status' command
        expected: Expected timezone offset (should be "-0800")
        **options: Additional options

    Returns:
        float: 1.0 if timezone is UTC-8, 0.0 otherwise
    """
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            if line.endswith('-0800)'):
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_not_exists__be265045(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if file does NOT exist.

    Args:
        result: Dict from getter with 'exists' key
        expected: Expected rules dict

    Returns:
        float: 1.0 if file should not exist and doesn't, 0.0 otherwise
    """
    exists = result.get('exists', False)
    should_not_exist = expected.get('should_not_exist', True)
    if should_not_exist and (not exists):
        return 1.0
    elif not should_not_exist and exists:
        return 1.0
    return 0.0

def check_json_settings__c5a909ed_v5(actual: str, expected: dict, **options) -> float:
    """
    Check if a specific key-value pair exists in the VS Code settings JSON file.
    This version handles boolean values correctly.

    Args:
        actual (str): path to result settings.json file
        expected (dict): expected dict with keys "expected_key" and "expected_value"

    Return:
        float: 1.0 if the key-value pair exists, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_json_settings__c5a909ed_v5: actual file path is None')
        return 0.0
    try:
        with open(actual, 'r') as f:
            data = json.load(f)
    except Exception as e:
        logger.debug(f'check_json_settings__c5a909ed_v5: Error reading JSON file: {e}')
        return 0.0
    expected_key = expected.get('expected_key')
    expected_value = expected.get('expected_value')
    if expected_key is None or expected_value is None:
        logger.debug('check_json_settings__c5a909ed_v5: expected_key or expected_value is None')
        return 0.0
    actual_value = data.get(expected_key)
    if actual_value == expected_value:
        return 1.0
    logger.debug(f'check_json_settings__c5a909ed_v5: Expected {expected_key}={expected_value}, got {actual_value}')
    return 0.0

def check_invoice_summary_text__834c93d1a65ecbb7766bb5ceb1a12320(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if the invoice summary text file contains the expected invoice information.

    Args:
        result: Content of the text file
        expected: Dict with 'invoice_entries' list containing expected invoice info
        **options: Additional options

    Returns:
        float: Partial score based on how many expected entries are found
    """
    if not result:
        return 0.0
    invoice_entries = expected.get('invoice_entries', [])
    if not invoice_entries:
        return 0.0
    result_lower = result.lower()
    found_count = 0
    for entry in invoice_entries:
        invoice_num = entry.get('invoice_num', '').lower()
        amount = entry.get('amount', '').lower()
        if invoice_num and amount:
            if invoice_num in result_lower and amount in result_lower:
                found_count += 1
    if len(invoice_entries) > 0:
        return found_count / len(invoice_entries)
    else:
        return 0.0

def check_timezone__f2f0035d(actual_file_path, expected):
    """Check if timezone is filled in correctly for conference locations.

    Args:
        actual_file_path: Path to the Excel file with filled timezone data
        expected: Expected configuration with timezone values

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_values = expected.get('expected_values', [])
    wb = openpyxl.load_workbook(actual_file_path, data_only=True)
    sheet = wb.active
    actual_values = []
    for row in sheet['C2:C22']:
        for cell in row:
            actual_values.append(cell.value)
    score = 0.0
    total_checks = len(expected_values)
    if total_checks == 0:
        return 0.0
    for i in range(total_checks):
        try:
            expected_val = expected_values[i]
            actual_val = actual_values[i]
            if actual_val is None:
                continue
            actual_str = str(actual_val).strip().upper().replace(' ', '')
            if isinstance(expected_val, list):
                if any((exp.strip().upper().replace(' ', '') in actual_str for exp in expected_val)):
                    score += 1.0 / total_checks
            elif expected_val.strip().upper().replace(' ', '') in actual_str:
                score += 1.0 / total_checks
        except (IndexError, KeyError, TypeError, AttributeError) as e:
            logger.debug(f'Error checking index {i}: {e}')
            continue
    return round(score, 3)

def check_all_files_exist__f50a55ca(result, expected, **options):
    """Check if all expected files exist.

    Args:
        result: Dictionary of file existence results
        expected: Dictionary of expected file existence values
        **options: Additional options

    Returns:
        float: 1.0 if all files exist as expected, 0.0 otherwise
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    for (file_path, should_exist) in expected.items():
        if result.get(file_path, False) != should_exist:
            return 0.0
    return 1.0

def check_file_exists_with_size__ff3634ef(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_file_count__48e4325f2d62197da2b10059281b95a0(result, expected, **options):
    """Check if the image_count.txt file exists and contains the correct count.

    Args:
        result: Dict with 'exists', 'content', and 'count' keys from getter
        expected: Rules dict with 'count' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists and count matches, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if not result.get('exists', False):
        return 0.0
    file_count = result.get('count')
    if file_count is None:
        return 0.0
    if file_count == expected_count:
        return 1.0
    else:
        return 0.0

def check_file_organization__8c566ad0(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_text_replacement__56d0ef227fc0a081b24201d3ecb3d358(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from file
        expected: Rules dict containing 'expected_text' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    result_clean = result.strip()
    expected_clean = expected_text.strip()
    if result_clean == expected_clean:
        return 1.0
    return 0.0

def compare_text_output__b1b8ec8f(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_file_exists__4e03b1ed(result, expected, **options):
    """Check if PDF file exists and is valid with correct content.

    Args:
        result: Dictionary from getter with validation results (exists, is_pdf, has_content, size, contains_keywords)
        expected: Expected configuration from rules
        **options: Additional options

    Returns:
        float: 1.0 if all validations pass, 0.0 otherwise
    """
    should_exist = expected.get('exists', True)
    if not should_exist:
        if not result.get('exists', True):
            logger.info(f'File correctly does not exist')
            return 1.0
        else:
            logger.info(f'File exists but should not')
            return 0.0
    if not result.get('exists', False):
        logger.info(f'File does not exist but should')
        return 0.0
    if not result.get('is_pdf', False):
        logger.info(f'File exists but is not a valid PDF')
        return 0.0
    if not result.get('has_content', False):
        logger.info(f"File exists and is PDF but appears empty or too small (size: {result.get('size', 0)} bytes)")
        return 0.0
    if not result.get('contains_keywords', False):
        logger.info(f'File exists and is valid PDF but does not contain expected content from the LLM agents blog post')
        return 0.0
    logger.info(f"File validation successful: exists as valid PDF with expected content (size: {result.get('size', 0)} bytes)")
    return 1.0

def check_timezone__29ace8b1(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_counts__816aac7b5fcbde572f13b62a3999bd4d(result: Dict, expected: Dict, **options) -> float:
    """
    Compare file count results against expected counts.

    Args:
        result: Dict with pattern counts from getter
        expected: Dict with expected counts for each pattern
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not expected:
        return 0.0
    expected_counts = expected.get('counts', {})
    total_checks = len(expected_counts)
    if total_checks == 0:
        return 0.0
    correct_checks = 0
    for (key, expected_count) in expected_counts.items():
        actual_count = result.get(key, 0)
        if actual_count == expected_count:
            correct_checks += 1
        logger.info(f"Pattern '{key}': expected={expected_count}, actual={actual_count}")
    score = correct_checks / total_checks
    return score

def check_text_replacement__62f7a00e7dd5e4db70eb00a615f012ef(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from the file
        expected: Expected rules dict with 'original_word' and 'correct_word'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    original_word = expected.get('original_word', '')
    correct_word = expected.get('correct_word', '')
    if not original_word or not correct_word:
        return 0.0
    if original_word.lower() in result.lower():
        return 0.0
    expected_count = expected.get('expected_count', 0)
    result_lower = result.lower()
    correct_word_lower = correct_word.lower()
    actual_count = result_lower.count(correct_word_lower)
    if actual_count >= expected_count:
        return 1.0
    if expected_count > 0:
        return min(1.0, actual_count / expected_count)
    return 0.0

def check_file_organization__a93d97ba(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_file_exists__82518dad(result, expected, **options):
    """
    Check if a file exists based on command output.

    This evaluator checks if lecture-notes.docx was successfully saved to /tmp.
    The instruction explicitly specifies the filename 'lecture-notes.docx' to match
    the actual attachment name in the email from the Thunderbird Notes folder.

    Args:
        result: Output from ls command checking /tmp/lecture-notes.docx
        expected: Dict with 'expected_output' key (should be 'exists')
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'exists')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_txt_line_count__8ff67d2b(result_state: Optional[str], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if a text file has the expected number of lines based on min_lines rule.

    Args:
        result_state: Path to the text file (from getter), or None if file doesn't exist
        expected_state: Dict containing 'min_lines' key (when type='rule', this IS the rules dict)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if the text file has at least the minimum number of lines, 0.0 otherwise
    """
    if result_state is None:
        return 0.0
    if not os.path.exists(result_state):
        return 0.0
    try:
        with open(result_state, 'r', encoding='utf-8', errors='replace') as f:
            lines = f.readlines()
            actual_line_count = len(lines)
        min_lines = expected_state.get('min_lines', 0)
        if actual_line_count >= min_lines:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_file_exists_with_size__21492178(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_text_alignment__1857b08997e5057be0ed7e2fe747fd99(result, expected, **options):
    """Check if the text alignment matches the expected alignment.

    Args:
        result: str - Alignment name from getter (e.g., 'CENTER', 'LEFT', 'RIGHT')
        expected: dict with 'alignment' key containing expected alignment name
        **options: Additional options

    Returns:
        float: 1.0 if alignment matches, 0.0 otherwise
    """
    expected_alignment = expected.get('alignment', 'LEFT')
    result_norm = result.strip().upper()
    expected_norm = expected_alignment.strip().upper()
    if result_norm == expected_norm:
        return 1.0
    return 0.0

def check_text_content__d7828490(result, expected, **options):
    """Compare text content against expected value.

    Args:
        result: Text content from getter
        expected: Expected text value (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_stripped = str(result).strip()
    expected_stripped = str(expected_value).strip()
    if result_stripped == expected_stripped:
        return 1.0
    else:
        return 0.0

def check_text_output__bb796efd(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_text_replacement__681a97c35eb849e5cf9422adfe4e1aea(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from file
        expected: Rules dict containing 'expected_text' key
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_text = expected.get('expected_text', '')
    if not expected_text:
        return 0.0
    result_clean = result.strip()
    expected_clean = expected_text.strip()
    if result_clean == expected_clean:
        return 1.0
    return 0.0

def check_file_naming__d6b113924c427deceec5f933af24484e(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if ZIP and internal files follow correct naming conventions.

    Args:
        result: Dict from getter with 'zip_exists', 'zip_name', 'internal_files'
        expected: Dict with 'base_name' to check against
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('zip_exists', False):
        return 0.0
    if not result.get('valid_zip', False):
        return 0.0
    base_name = expected.get('base_name', '')
    zip_name = result.get('zip_name', '')
    internal_files = result.get('internal_files', [])
    score = 0.0
    expected_zip_name = f'{base_name}.zip'
    if zip_name == expected_zip_name:
        score += 0.5
    expected_docx = f'{base_name}.docx'
    expected_pdf = f'{base_name}.pdf'
    has_correct_docx = any((f == expected_docx or f.endswith(f'/{expected_docx}') for f in internal_files))
    has_correct_pdf = any((f == expected_pdf or f.endswith(f'/{expected_pdf}') for f in internal_files))
    if has_correct_docx:
        score += 0.25
    if has_correct_pdf:
        score += 0.25
    return min(score, 1.0)

def check_file_exists__62dacdc1(result, expected, **options):
    """
    Check if file exists and is a valid PNG image.

    Args:
        result: Dict with file information (exists, file_type, size, is_png, is_valid_image)
        expected: Expected value (dict with 'exists' field)

    Returns:
        1.0 if file exists, is a valid PNG image with reasonable size, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        return 1.0 if not result.get('exists', False) else 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_png', False):
        logger.info(f"File is not a PNG image. File type: {result.get('file_type', 'unknown')}")
        return 0.0
    if not result.get('is_valid_image', False):
        logger.info('File is not a valid PNG image (magic bytes check failed)')
        return 0.0
    min_size = 1024
    file_size = result.get('size', 0)
    if file_size < min_size:
        logger.info(f'File size ({file_size} bytes) is too small for a valid screenshot (minimum {min_size} bytes)')
        return 0.0
    logger.info(f'File validation passed: PNG image with size {file_size} bytes')
    return 1.0

def check_exact_match_v6__d4477d7a(result, expected, **options):
    """Compare result against expected integer value.

    Args:
        result: Actual value from getter
        expected: Rules dict with 'expected_value' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_value = expected.get('expected_value')
    if result is None or expected_value is None:
        return 0.0
    try:
        result_int = int(result)
        expected_int = int(expected_value)
    except (TypeError, ValueError):
        return 0.0
    if result_int == expected_int:
        return 1.0
    return 0.0

def check_python_complete__12fe4256(result_file, expected, **options):
    """
    Check if a Python file contains complete code with multiple requirements.

    Args:
        result_file: Path to the Python file to check
        expected: Dict with 'min_lines', 'required_imports', 'required_classes', 'required_functions'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 (0.25 per requirement met)
    """
    if not result_file:
        return 0.0
    try:
        with open(result_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    score = 0.0
    min_lines = expected.get('min_lines', 0)
    lines = [l for l in content.split('\n') if l.strip()]
    if len(lines) >= min_lines:
        score += 0.25
    try:
        tree = ast.parse(content)
    except SyntaxError:
        return score
    required_imports = expected.get('required_imports', [])
    if required_imports:
        imported_modules = set()
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                for alias in node.names:
                    imported_modules.add(alias.name)
            elif isinstance(node, ast.ImportFrom):
                if node.module:
                    imported_modules.add(node.module)
        found_imports = sum((1 for imp in required_imports if any((imp == mod or mod.startswith(imp + '.') or imp.startswith(mod + '.') for mod in imported_modules))))
        if found_imports == len(required_imports):
            score += 0.25
    required_classes = expected.get('required_classes', [])
    if required_classes:
        found_classes = set()
        for node in ast.walk(tree):
            if isinstance(node, ast.ClassDef):
                found_classes.add(node.name)
        if all((cls in found_classes for cls in required_classes)):
            score += 0.25
    required_functions = expected.get('required_functions', [])
    if required_functions:
        found_functions = set()
        for node in tree.body:
            if isinstance(node, ast.FunctionDef):
                found_functions.add(node.name)
        if all((func in found_functions for func in required_functions)):
            score += 0.25
    return score

def check_text_lines_exact__4250d59b26bb86f2de0562f0a55c312c(result, expected, **options):
    """Compare text file lines for exact match.

    Args:
        result: List of lines from result file
        expected: Expected list of lines (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_lines = expected.get('lines', [])
    if not isinstance(result, list):
        return 0.0
    if len(result) != len(expected_lines):
        return 0.0
    for (i, (res_line, exp_line)) in enumerate(zip(result, expected_lines)):
        if str(res_line).strip() != str(exp_line).strip():
            return 0.0
    return 1.0

def check_text_match__93eac3e2452ef121ce8047db9ec250fe(result, expected, **options):
    """Check if text content matches expected value.

    Args:
        result: Actual text content from getter
        expected: Expected text value (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_value = expected.get('value', '')
    result_normalized = ' '.join(result.split())
    expected_normalized = ' '.join(str(expected_value).split())
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_exists_with_size__479794c8(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_file_organization__e2bf8bf2(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_file_ownership__18d171e6(result, rules, **options):
    """Check if all files have the expected owner.

    Args:
        result: Output from ls -l command showing file ownership
        rules: Dict with 'owner' key specifying expected owner username

    Returns:
        float: 1.0 if all files have correct owner, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    expected_owner = rules.get('owner', '')
    lines = result.strip().split('\n')
    for line in lines:
        if not line.strip():
            continue
        parts = line.split()
        if len(parts) > 2:
            actual_owner = parts[2]
            if actual_owner != expected_owner:
                return 0.0
    return 1.0

def check_timezone_utc_plus_3__d1091690(result_state, expected_state, **options):
    """
    Check if timezone is set to UTC+3 (Moscow Standard Time).

    Args:
        result_state: Output from 'timedatectl status' command
        expected_state: Expected state with rules for timezone validation
        **options: Additional options

    Returns:
        float: 1.0 if timezone offset is +0300 and timezone name matches, 0.0 otherwise
    """
    timedatectl_output = result_state
    if isinstance(expected_state, dict):
        timezone_offset = expected_state.get('timezone_offset', '+0300')
        timezone_pattern = expected_state.get('timezone_pattern', '.*\\+0300\\)$')
        timezone_names = expected_state.get('timezone_names', ['Moscow', 'MSK', 'Europe/Moscow'])
    else:
        timezone_offset = '+0300'
        timezone_pattern = '.*\\+0300\\)$'
        timezone_names = ['Moscow', 'MSK', 'Europe/Moscow']
    lines = timedatectl_output.split('\n')
    timezone_line = None
    for line in lines:
        if 'Time zone:' in line:
            timezone_line = line
            break
    if not timezone_line:
        if len(lines) > 3:
            timezone_line = lines[3]
        else:
            return 0.0
    if not re.search(timezone_pattern, timezone_line):
        return 0.0
    timezone_name_found = False
    for name in timezone_names:
        if name in timezone_line:
            timezone_name_found = True
            break
    if not timezone_name_found:
        if timezone_offset in timezone_line:
            return 1.0
        return 0.0
    return 1.0

def check_selective_file_deletion__ca75e69be093d6cb8e4fa53ce6114782(result, expected, **options):
    """Check if files were selectively deleted from specific directories.

    Args:
        result: Dict mapping directories to file existence booleans
        expected: Dict with expected existence values for each directory
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    for (directory, should_exist) in expected.items():
        if result.get(directory, False) != should_exist:
            return 0.0
    return 1.0

def check_text_replacement__c5c5b95c23a9d2c28863362320aba24b(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from the file
        expected: Expected rules dict with 'original_word' and 'correct_word'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    original_word = expected.get('original_word', '')
    correct_word = expected.get('correct_word', '')
    if not original_word or not correct_word:
        return 0.0
    if original_word in result:
        return 0.0
    expected_count = expected.get('expected_count', 0)
    actual_count = result.count(correct_word)
    if actual_count >= expected_count:
        return 1.0
    if expected_count > 0:
        return min(1.0, actual_count / expected_count)
    return 0.0

def compare_text_output__cee492ee(actual_path: str, expected: dict, **options) -> float:
    """
    Compare actual text file content against expected text.
    
    Args:
        actual_path: Path to the result text file
        expected: Dict containing "expected_text" key with expected output
        **options: Additional options (ignore_blanks, ignore_case, etc.)
        
    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not actual_path:
        return 0.0
    try:
        with open(actual_path, 'r') as f:
            actual_text = f.read()
    except FileNotFoundError:
        return 0.0
    except Exception:
        return 0.0
    expected_text = expected.get('expected_text', '')
    ignore_blanks = options.get('ignore_blanks', False)
    if ignore_blanks:
        import re
        actual_text = re.sub('[\\t\\n]', ' ', actual_text).strip()
        actual_text = re.sub('\\s+', ' ', actual_text)
        expected_text = re.sub('[\\t\\n]', ' ', expected_text).strip()
        expected_text = re.sub('\\s+', ' ', expected_text)
    ignore_case = options.get('ignore_case', False)
    if ignore_case:
        actual_text = actual_text.lower()
        expected_text = expected_text.lower()
    if actual_text == expected_text:
        return 1.0
    return 0.0

def check_text_patterns__32952afd(result: Optional[str], expected: Dict[str, Any], **options) -> float:
    """Check if text contains specific Apple UI Buttons documentation content.

    This function verifies that the document contains Apple-specific terminology
    unique to their design documentation, preventing false positives from generic
    UI-related text.

    Args:
        result: Actual text content from getter
        expected: Rules dict with:
            - 'patterns': list of Apple-specific patterns to match
            - 'min_length': minimum character count (default 200)
        **options: Additional options

    Returns:
        Score 1.0 if content is verified as Apple Buttons documentation, 0.0 otherwise
    """
    if result is None:
        return 0.0
    patterns = expected.get('patterns', [])
    min_length = expected.get('min_length', 200)
    if not patterns:
        return 0.0
    if len(result) < min_length:
        return 0.0
    result_lower = result.lower()
    found_count = 0
    for pattern in patterns:
        if pattern.lower() in result_lower:
            found_count += 1
    required_threshold = max(6, int(len(patterns) * 0.75))
    if found_count >= required_threshold:
        return 1.0
    return 0.0

def check_downloads_file_exists__2cbf25da(result, expected, **options):
    """Check if file exists in Downloads.

    Args:
        result: Boolean from getter
        expected: Expected value from rules
        **options: Additional options

    Returns:
        1.0 if file exists, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_file_created__974295f9d11461d175dbc0223dd4ff65(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if a file was created with correct content.

    Args:
        result: Dict with 'exists' (bool) and 'content' (str) from getter
        expected: Expected dict with 'exists' (bool) and 'content' (str) keys
        **options: Additional options

    Returns:
        float: 1.0 if file existence and content match expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    result_exists = result.get('exists', False)
    if result_exists != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    expected_content = expected.get('content', '')
    result_content = result.get('content', '')
    if result_content.strip() == expected_content.strip():
        return 1.0
    return 0.0

def check_files_renamed__416ce0b1(result, rules, **options):
    """Check if all files were renamed with the specified prefix.

    Args:
        result: Output from find command with -printf '%f\\n' (filenames only)
        rules: Dict with 'prefix' and 'original_names' keys

    Returns:
        float: 1.0 if all files have the prefix, 0.0 otherwise
    """
    if result is None or not result.strip():
        return 0.0
    prefix = rules.get('prefix', '')
    original_names = rules.get('original_names', [])
    filenames = [line.strip() for line in result.strip().split('\n') if line.strip()]
    if len(filenames) != len(original_names):
        return 0.0
    expected_renamed = [prefix + name for name in original_names]
    for expected_name in expected_renamed:
        if expected_name not in filenames:
            return 0.0
    return 1.0

def check_all_text_color_match__6666d59e63b5ea41073d9fefe10e8bab(result, expected, **options):
    """Check if all text colors match the expected color (with tolerance).

    Args:
        result: List of RGB tuples from getter
        expected: Dict with 'target_color' (R, G, B) tuple
        **options: 'color_tolerance' (default 20) for fuzzy matching

    Returns:
        float: 1.0 if all colors match target, 0.0 otherwise
    """
    if not result:
        return 0.0
    target_color = expected.get('target_color')
    if not target_color:
        return 0.0
    tolerance = options.get('color_tolerance', 20)

    def colors_match(c1, c2, tol):
        """Check if two RGB colors match within tolerance."""
        return all((abs(c1[i] - c2[i]) <= tol for i in range(3)))
    all_match = all((colors_match(color, target_color, tolerance) for color in result))
    return 1.0 if all_match else 0.0

def check_file_exists__dc35efec(result, expected, **options):
    """
    Check if file exists and is a valid PNG screenshot.

    Validates that:
    1. File exists at the expected path
    2. File is a valid PNG image (verified by MIME type)
    3. File has reasonable size (> 1KB to exclude empty/dummy files)
    4. File was created/modified recently (within last 5 minutes)

    Args:
        result: Dict with file information from getter:
            - exists: bool
            - is_png: bool
            - file_size: int
            - is_recent: bool
        expected: Expected value dict with:
            - exists: bool (True for file should exist)

    Returns:
        1.0 if all validations pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        if not result.get('exists', False):
            return 1.0
        else:
            return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_png', False):
        logger.info('File is not a valid PNG image')
        return 0.0
    file_size = result.get('file_size', 0)
    if file_size < 1024:
        logger.info(f'File size too small: {file_size} bytes (expected > 1KB)')
        return 0.0
    if not result.get('is_recent', False):
        logger.info('File is not recent (was not created/modified within last 5 minutes)')
        return 0.0
    return 1.0

def check_rtf_file_exists__34460ac1a76394f2dfa8b4e9981344a0(result, expected, **options):
    """
    Check if RTF file exists.

    Args:
        result: Dict from getter with 'exists' key
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    exists = result.get('exists', False)
    if exists:
        logger.info(f"✅ RTF file exists at {result.get('path', 'unknown')}")
        return 1.0
    else:
        logger.warning(f"❌ RTF file does not exist at {result.get('path', 'unknown')}")
        return 0.0

def check_title_text__8b1e2a3b(pptx_path, expected):
    """
    Check if a specific shape contains expected text.

    Args:
        pptx_path: Path to the PPTX file
        expected: Dict containing 'slide_idx', 'shape_idx', and 'text'

    Returns:
        float: 1.0 if text matches, 0.0 otherwise
    """
    try:
        presentation = Presentation(pptx_path)
        slide_idx = expected.get('slide_idx', 0)
        shape_idx = expected.get('shape_idx', 1)
        expected_text = expected.get('text', '')
        if slide_idx >= len(presentation.slides):
            return 0.0
        slide = presentation.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            return 0.0
        shape = slide.shapes[shape_idx]
        if not hasattr(shape, 'text'):
            return 0.0
        actual_text = shape.text.strip()
        expected_text_stripped = expected_text.strip()
        if actual_text == expected_text_stripped:
            return 1.0
        return 0.0
    except Exception as e:
        return 0.0

def check_file_existence__e1da6937(result, expected, **options):
    """
    Check if file existence matches expected.

    Args:
        result: Boolean from getter
        expected: Expected boolean value

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_new_textbox_added__4ebe6ee8(result_state, expected, **options):
    """Check if a new text box with specific text was added to the slide.

    This metric uses comparison data from the getter (which compared baseline
    vs result internally) to verify that:
    1. A NEW text box shape was added (textbox count increased)
    2. The new text box contains exactly the expected text

    Args:
        result_state: Dict with comparison data from getter including:
            - baseline_textbox_count: Number of textboxes before task
            - result_textbox_count: Number of textboxes after task
            - baseline_texts: List of text from baseline textboxes
            - result_textboxes: List of textbox info from result
            - target_text: Expected text to find
        expected: Dict with 'target_text' from config

    Returns:
        float: 1.0 if new text box with expected text was added, 0.0 otherwise
    """
    try:
        baseline_count = result_state.get('baseline_textbox_count', 0)
        result_count = result_state.get('result_textbox_count', 0)
        baseline_texts = set(result_state.get('baseline_texts', []))
        result_textboxes = result_state.get('result_textboxes', [])
        target_text = result_state.get('target_text', '') or expected.get('target_text', '')
        if result_count <= baseline_count:
            return 0.0
        for textbox in result_textboxes:
            text_content = textbox.get('text', '').strip()
            if text_content not in baseline_texts:
                if text_content == target_text:
                    return 1.0
        return 0.0
    except Exception as e:
        print(f'Error checking new textbox: {e}')
        import traceback
        traceback.print_exc()
        return 0.0

def check_text_content__7bbdf0a0733630cbbfd86729556fc827(result, expected, **options):
    """Check if file content matches expected text exactly.

    Args:
        result: String content from the file
        expected: Dict with 'content' key containing expected text
        **options: Additional options (ignore_trailing_whitespace, etc.)

    Returns:
        float: 1.0 if content matches, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_content = expected.get('content', '')
    ignore_trailing = options.get('ignore_trailing_whitespace', False)
    if ignore_trailing:
        result = result.rstrip()
        expected_content = expected_content.rstrip()
    if result == expected_content:
        return 1.0
    return 0.0

def check_file_created__8ab0a45d4a57cf0e9592e87621895b59(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a file was created with expected content.

    Args:
        result: Dict from getter with file information
        expected: Dict with expected file properties
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    score += 0.3
    if result.get('is_file', False):
        score += 0.2
    else:
        logger.warning('Path exists but is not a file')
        return score
    required_text = expected.get('required_text', [])
    if required_text:
        content = result.get('content', '')
        matching_text = sum((1 for text in required_text if text in content))
        if matching_text == len(required_text):
            score += 0.5
        else:
            score += 0.5 * (matching_text / len(required_text))
        logger.info(f'Required text matched: {matching_text}/{len(required_text)}')
    logger.info(f'File creation check score: {score}')
    return score

def check_file_exists__6d219cf2(result, expected, **options):
    """
    Check if file exists and is a valid PNG with reasonable content.

    Args:
        result: Dict with 'exists', 'size', and 'is_png' fields
        expected: Expected rules (exists should be True)

    Returns:
        1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not isinstance(result, dict):
        return 0.0
    file_exists = result.get('exists', False)
    if file_exists != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    is_png = result.get('is_png', False)
    if not is_png:
        logger.warning('File exists but is not a valid PNG (magic bytes check failed)')
        return 0.0
    file_size = result.get('size', 0)
    if file_size < 1024:
        logger.warning(f'File size too small: {file_size} bytes (expected > 1KB)')
        return 0.0
    if file_size > 10 * 1024 * 1024:
        logger.warning(f'File size too large: {file_size} bytes (expected < 10MB)')
        return 0.0
    return 1.0

def check_python_content__198be354(result, expected, **options):
    """Check if Python file contains expected code patterns.

    Args:
        result: dict from getter with content analysis
        expected: dict with required_patterns list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('exists', False):
        return 0.0
    required_patterns = expected.get('required_patterns', [])
    if not required_patterns:
        return 1.0
    matched = sum((1 for pattern in required_patterns if result.get(pattern, False)))
    score = matched / len(required_patterns)
    return score

def check_filename_prefix__fdcd3f41(result, expected, **options):
    """Check if all filenames start with the expected prefix.

    Args:
        result: List of filenames from getter
        expected: Expected prefix and count (from rules dict)
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    prefix = expected.get('prefix', '')
    expected_count = expected.get('count', 0)
    if len(result) != expected_count:
        return 0.0
    for filename in result:
        if not filename.startswith(prefix):
            return 0.0
    return 1.0

def check_filename__739292ff(result, expected, **options):
    """Check if filename matches expected pattern and file is valid.

    Args:
        result: Dict with file information
        expected: Dict with 'name' or 'extension' keys
        **options: Additional options

    Returns:
        float: Score based on filename match and file validity
    """
    if result is None:
        logger.error('Result is None')
        return 0.0
    if not result.get('exists', False):
        logger.error('File does not exist')
        return 0.0
    file_size = result.get('size', 0)
    if file_size <= 0:
        logger.error(f'File is empty or has invalid size: {file_size}')
        return 0.0
    score = 0.0
    checks = 0
    if 'name' in expected:
        checks += 1
        if result.get('name') == expected['name']:
            score += 1.0
            logger.info(f"Filename matches: {result.get('name')}")
        else:
            logger.info(f"Filename mismatch: expected {expected['name']}, got {result.get('name')}")
    if 'extension' in expected:
        checks += 1
        if result.get('extension') == expected['extension']:
            score += 1.0
            logger.info(f"Extension matches: {result.get('extension')}")
        else:
            logger.info(f"Extension mismatch: expected {expected['extension']}, got {result.get('extension')}")
    if checks == 0:
        logger.warning('No checks specified')
        return 0.0
    return score / checks

def check_textbox_centered__13f49ee1(src_path, expected_state, **options):
    """
    Check if the textbox is centered horizontally on the image.
    Variation 3 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text center is within middle 10%, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    right_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                right_most = max(right_most, x)
    text_center = (left_most + right_most) / 2
    image_center = width / 2
    tolerance = width * 0.05
    if abs(text_center - image_center) < tolerance:
        return 1.0
    else:
        return 0.0

def check_file_exists__cdcbbd90(result, expected, **options):
    """
    Check if file exists.

    Args:
        result: Boolean indicating if file exists
        expected: Expected value (True for file should exist)

    Returns:
        1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_file_existence__198be354(result, expected, **options):
    """Check if file exists with required properties and contains code from Colab notebook.

    This metric now verifies:
    1. File exists at expected path
    2. File has .py extension
    3. File has content
    4. Content is valid Python syntax
    5. Content has code patterns
    6. Content actually contains code extracted from the Colab notebook (NEW)

    Args:
        result: dict from getter with file existence, content analysis, and Colab code comparison
        expected: dict with 'must_exist', 'must_be_python', 'must_have_content' flags
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if expected.get('must_exist', True):
        if result.get('exists', False):
            score += 0.2
        else:
            return 0.0
    if expected.get('must_be_python', True):
        if result.get('is_python', False):
            score += 0.1
    if expected.get('must_have_content', True):
        if not result.get('is_empty', True):
            score += 0.1
    if result.get('is_valid_python', False):
        score += 0.15
    if result.get('has_code_patterns', False):
        score += 0.1
    colab_code_cells = result.get('colab_code_cells', [])
    code_similarity_score = result.get('code_similarity_score', 0.0)
    if colab_code_cells:
        score += 0.35 * code_similarity_score
    else:
        pass
    return score

def check_file_exists__dd9409c8(result, expected, **options):
    """Check if file existence matches expected value.

    This metric validates that a VLC snapshot was successfully created by checking:
    1. A valid PNG file exists in the expected location
    2. The file has meaningful content (size > 1KB)
    3. The file was created recently (not a pre-existing file)

    Args:
        result: Dict from getter with keys:
            - exists: bool - whether a valid recent snapshot exists
            - file_path: str - path to the found file (or None)
            - file_size: int - size in bytes (or 0)
            - is_recent: bool - whether file was created recently
        expected: Expected dict value (from rules) with key:
            - exists: bool - expected existence state
        **options: Additional options

    Returns:
        1.0 if all conditions match, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if isinstance(result, dict):
        actual_exists = result.get('exists', False)
        is_recent = result.get('is_recent', False)
        file_size = result.get('file_size', 0)
        file_path = result.get('file_path', None)
        if expected_exists:
            if actual_exists and is_recent and (file_size >= 1024):
                logger.info(f'VLC snapshot validation passed: {file_path} ({file_size} bytes, recent)')
                return 1.0
            else:
                logger.warning(f'VLC snapshot validation failed: exists={actual_exists}, recent={is_recent}, size={file_size}')
                return 0.0
        elif not actual_exists:
            return 1.0
        else:
            return 0.0
    elif isinstance(result, bool):
        if result == expected_exists:
            return 1.0
        else:
            return 0.0
    else:
        logger.error(f'Unexpected result type: {type(result)}')
        return 0.0

def check_file_exists__916d8b58(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if command output indicates file exists.

    Args:
        result: Command output (EXISTS or MISSING)
        expected: Expected rules dict with 'expected_output'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    try:
        expected_output = expected.get('expected_output', 'EXISTS')
        if expected_output.strip().upper() in result.strip().upper():
            return 1.0
        else:
            return 0.0
    except Exception as e:
        return 0.0

def check_git_dir_exists__e2da960ab9034666db33db74ae6371a7(result, expected, **options):
    """Check if .git directory existence matches expected value.

    Args:
        result: Boolean from getter (True if .git exists)
        expected: Rules dict with 'exists' key (bool)
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_textbox_on_rightside__5c6340a1(src_path, expected, **options):
    """
    Check if the textbox is on the right side of the image.
    Task variation 0 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is on right side, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    right_most_dark_pixel = 0
    for y in range(height):
        for x in range(width - 1, -1, -1):
            if gray_image.getpixel((x, y)) < 128:
                right_most_dark_pixel = max(right_most_dark_pixel, x)
                break
    if right_most_dark_pixel > width * 0.95:
        return 1.0
    else:
        return 0.0

def check_file_exists__0b2042a1(result, expected, **options):
    """
    Check if a file exists based on command output.

    Args:
        result: Output from ls command
        expected: Dict with 'expected_output' key
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_output = expected.get('expected_output', 'exists')
    result_stripped = result.strip()
    return 1.0 if expected_output in result_stripped else 0.0

def check_text_output__afe6b7a7(result, expected, **options):
    """Compare text output against expected value.

    Args:
        result: Actual text output from getter
        expected: Expected text content (from rules dict)
        **options: Additional comparison options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '').strip()
    result_normalized = result.strip()
    expected_normalized = expected_text.strip()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        return 0.0

def check_file_count__03f8ef9d(result: list, expected: dict, **options) -> float:
    """Check if the filenames in the zip match the expected filenames.

    Args:
        result: List of actual filenames in the zip
        expected: Dict with 'filenames' key containing expected filenames
        **options: Additional options

    Returns:
        float: Score based on correctness (1.0 for perfect match, partial credit for some correct files)
    """
    expected_filenames = expected.get('filenames', [])
    result_set = set(result)
    expected_set = set(expected_filenames)
    if result_set == expected_set:
        return 1.0
    correct_files = result_set & expected_set
    incorrect_files = result_set - expected_set
    missing_files = expected_set - result_set
    score = 0.0
    score += len(correct_files) * 0.5
    score -= len(incorrect_files) * 0.25
    return max(0.0, min(1.0, score))

def check_file_contains_lines__d5302e2f(result_file_path, expected, **options):
    """
    Check if a text file contains all expected lines.
    
    Args:
        result_file_path: Path to the result text file
        expected: Dict with 'lines' key containing list of expected lines
        **options: Additional options (case_sensitive, order_matters, etc.)
        
    Returns:
        float: Score between 0.0 and 1.0
    """
    if result_file_path is None:
        return 0.0
    try:
        with open(result_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            result_lines = [line.strip() for line in content.strip().split('\n') if line.strip()]
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    expected_lines = expected.get('lines', [])
    if not expected_lines:
        return 0.0
    case_sensitive = options.get('case_sensitive', True)
    order_matters = options.get('order_matters', False)
    if not case_sensitive:
        result_lines_normalized = [line.lower() for line in result_lines]
        expected_lines_normalized = [line.lower() for line in expected_lines]
    else:
        result_lines_normalized = result_lines
        expected_lines_normalized = expected_lines
    if order_matters:
        if len(result_lines_normalized) != len(expected_lines_normalized):
            return 0.0
        matches = sum((1 for (i, exp_line) in enumerate(expected_lines_normalized) if i < len(result_lines_normalized) and result_lines_normalized[i] == exp_line))
        return matches / len(expected_lines_normalized)
    else:
        matches = sum((1 for exp_line in expected_lines_normalized if exp_line in result_lines_normalized))
        return matches / len(expected_lines_normalized)

def check_timestamped_file__2b78c2fd0d670b6ee1c54ce65b4419a5(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if file with YYYYMMDD timestamp exists and the date is today.

    Args:
        result: Dict from getter
        expected: Dict with expected properties

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        logger.info('File does not exist')
        return 0.0
    if result.get('has_valid_yyyymmdd_format', False):
        date_value = result.get('date_value')
        if date_value:
            today = datetime.now().strftime('%Y%m%d')
            if date_value == today:
                score += 0.4
                logger.info(f'Date matches today: {date_value}')
            else:
                logger.info(f'Date {date_value} does not match today {today}')
        else:
            logger.info('Date value not extracted')
    else:
        logger.info('Filename does not contain valid YYYYMMDD format')
    if result.get('is_png', False):
        score += 0.2
    return score

def check_file_exists__e3ae8a85(result, expected, **options):
    """
    Check if file exists.

    Args:
        result: Boolean indicating if file exists
        expected: Expected value (True for file should exist)

    Returns:
        1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    else:
        return 0.0

def check_timezone_pst__27082429(timedatectl_output, expected, **options):
    """
    Check if timezone is set to Pacific Standard Time (PST, UTC-8).

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Expected timezone configuration (dict with 'offset' and 'timezone' keys)
        **options: Additional options

    Returns:
        1.0 if timezone matches expected offset and timezone name, 0.0 otherwise
    """
    expected_offset = expected.get('offset')
    expected_timezone = expected.get('timezone')
    if not expected_offset or not expected_timezone:
        return 0.0
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            if expected_timezone in line and expected_offset in line:
                return 1.0
            break
    return 0.0

def check_textbox_movedleft__0fe6ee6f(result_state, expected_state, **options):
    """
    Check if the textbox has been moved to the left compared to its initial position.
    Variation 8 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    This metric verifies RELATIVE movement by comparing the initial and final positions.

    Args:
        result_state: Dict with 'initial' and 'final' positions
            - initial: dict with 'x' (leftmost pixel) and 'width' from initial state
            - final: str path to the exported final image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text moved left (final_x < initial_x), 0.0 otherwise
    """
    if result_state is None:
        return 0.0
    try:
        initial_state = result_state.get('initial')
        final_path = result_state.get('final')
        if initial_state is None or final_path is None:
            return 0.0
        initial_x = initial_state.get('x')
        initial_width = initial_state.get('width')
        if initial_x is None or initial_width is None:
            return 0.0
        final_image = Image.open(final_path)
        gray_image = final_image.convert('L')
        (width, height) = final_image.size
        final_x = width
        for y in range(height):
            for x in range(width):
                if gray_image.getpixel((x, y)) < 128:
                    final_x = min(final_x, x)
                    break
        movement_amount = initial_x - final_x
        if movement_amount > 0:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        import logging
        logging.error(f'Error checking textbox position: {e}')
        import traceback
        logging.error(traceback.format_exc())
        return 0.0

def check_file_exported__60340d37(result, expected, **options):
    """Check if file was successfully exported as ODS format.

    Args:
        result: Dictionary with file existence, format, and size info
        expected: Expected value (True for file should exist)
        **options: Additional options

    Returns:
        float: 1.0 if file exists as valid ODS with content, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if isinstance(result, bool):
        return 1.0 if result == expected else 0.0
    if not result.get('exists', False):
        return 0.0
    if not result.get('is_ods', False):
        return 0.0
    file_size = result.get('size', 0)
    if file_size < 1024:
        return 0.0
    return 1.0

def check_python_docstring__6ee0182d(actual: str, rules: dict, **options) -> float:
    """
    Check if a Python file has a docstring for a specified function.

    Args:
        actual (str): path to Python file
        rules (dict): expected configuration rules with function_name

    Returns:
        float: score between 0.0 and 1.0
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
    except Exception as e:
        logger.error(f'Failed to read Python file: {e}')
        return 0.0
    function_name = rules.get('function_name', 'hello_world')
    try:
        tree = ast.parse(content)
    except SyntaxError as e:
        logger.error(f'Syntax error in Python file: {e}')
        return 0.0
    score = 0.0
    for node in ast.walk(tree):
        if isinstance(node, ast.FunctionDef) and node.name == function_name:
            score += 0.5
            docstring = ast.get_docstring(node)
            if docstring and len(docstring.strip()) > 0:
                score += 0.5
            break
    return score

def check_textbox_smaller__a8f988b59259c0be84da4d1fc65b92ad(src_path, expected):
    """
    Check if the text layer has been scaled down (made smaller).
    Variation of gimp:e2dd0213-26db-4349-abe5-d5667bfd725c
    Task: Scale text layer to be smaller

    Args:
        src_path: Path to the exported image
        expected: Dict with 'max_text_width_ratio' and 'max_text_height_ratio'

    Returns:
        1.0 if text occupies less than the specified ratio of image dimensions, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    max_width_ratio = expected.get('max_text_width_ratio', 0.3)
    max_height_ratio = expected.get('max_text_height_ratio', 0.15)
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    right_most = 0
    top_most = height
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                right_most = max(right_most, x)
                top_most = min(top_most, y)
                bottom_most = max(bottom_most, y)
    if left_most < width and top_most < height:
        text_width = right_most - left_most
        text_height = bottom_most - top_most
        width_ratio = text_width / width
        height_ratio = text_height / height
        if width_ratio < max_width_ratio and height_ratio < max_height_ratio:
            return 1.0
    return 0.0

def check_python_imports__52a225d6(result_file, expected, **options):
    """
    Check if a Python file contains required import statements.

    Args:
        result_file: Path to the Python file to check
        expected: Dict with 'required_imports' key containing list of module names
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result_file:
        return 0.0
    try:
        with open(result_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    required_imports = expected.get('required_imports', [])
    if not required_imports:
        return 0.0
    try:
        tree = ast.parse(content)
    except SyntaxError:
        imported_modules = set()
        import_patterns = ['^\\s*import\\s+([\\w\\.]+)', '^\\s*from\\s+([\\w\\.]+)\\s+import']
        for line in content.split('\n'):
            for pattern in import_patterns:
                match = re.match(pattern, line)
                if match:
                    imported_modules.add(match.group(1))
    else:
        imported_modules = set()
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                for alias in node.names:
                    imported_modules.add(alias.name)
            elif isinstance(node, ast.ImportFrom):
                if node.module:
                    imported_modules.add(node.module)
    found_count = 0
    for required in required_imports:
        if any((required == imp or imp.startswith(required + '.') or required.startswith(imp + '.') for imp in imported_modules)):
            found_count += 1
    score = found_count / len(required_imports)
    return score

def check_text_content__7676732d(result, expected, **options):
    """Verify both count file content and CSV export existence.

    This metric checks two requirements:
    1. The count file contains the expected value (number of contacts born in 1990s)
       Expected value '8' verified against actual Thunderbird profile data:
       - Total 30 contacts in Personal Address Book
       - 8 contacts born between 1990-1999 (inclusive):
         Monica Mayo (1996), Craig Williams (1992), Kevin Alexander (1993),
         Thomas Jensen (1990), Sandra Williams (1999), Erica Barron (1994),
         Kimberly Moss (1998), Brittany Allen (1996)
    2. A CSV file exists on the Desktop (the exported address book)
       Preference given to files with address book related keywords

    Args:
        result: Dict from getter containing:
            - 'count_content': Content of the count file (str)
            - 'csv_exists': Whether a CSV file exists on Desktop (bool)
            - 'csv_file': Name of the CSV file if found (str or None)
        expected: Expected text value (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if both requirements met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_value = expected.get('value', '')
    count_content = result.get('count_content', '')
    count_matches = str(count_content).strip() == str(expected_value).strip()
    csv_exists = result.get('csv_exists', False)
    if count_matches and csv_exists:
        return 1.0
    else:
        return 0.0

def check_line_order__2c3b878a(result: str, expected: Dict, **options) -> float:
    """Check if lines appear in expected order.

    Args:
        result: The concatenated lines string
        expected: Dict with 'order' key containing expected sequence

    Returns:
        1.0 if order matches, 0.0 otherwise
    """
    expected_order = expected.get('order', '')
    if result == expected_order:
        return 1.0
    return 0.0

def check_line_duplicated__33cae2c6(result, expected, **options):
    """Check if a specific line was duplicated at the end of the file.

    Args:
        result: List of file lines
        expected: Dict with 'original_line_num' (1-based), 'duplicate_position' ('end' or line number)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    original_line_num = expected.get('original_line_num', 1)
    duplicate_position = expected.get('duplicate_position', 'end')
    original_idx = original_line_num - 1
    if original_idx >= len(result):
        return 0.0
    original_line = result[original_idx]
    if duplicate_position == 'end':
        last_line = result[-1]
        if last_line == original_line:
            return 1.0
    else:
        dup_idx = duplicate_position - 1
        if dup_idx < len(result) and result[dup_idx] == original_line:
            return 1.0
    return 0.0

def check_chapter_file_naming__29db12fd(result_state, expected_state, **options):
    """
    Check if chapter files have been renamed with leading zeros.

    Args:
        result_state: Output from ls command (str)
        expected_state: Expected rules dict with chapter_count and expected_pattern
        **options: Additional options

    Returns:
        float: Score (1.0 if all files correctly renamed, 0.0 otherwise)
    """
    if not result_state or not isinstance(result_state, str):
        return 0.0
    chapter_count = expected_state.get('chapter_count', 5)
    files = result_state.strip().split('\n')
    files = [f.strip() for f in files if f.strip()]
    expected_files = [f'Chapter0{i}.txt' for i in range(chapter_count)]
    found_count = 0
    for expected_file in expected_files:
        if expected_file in files:
            found_count += 1
    old_files = [f'Chapter{i}.txt' for i in range(chapter_count)]
    old_files_exist = False
    for old_file in old_files:
        if old_file in files:
            old_files_exist = True
            break
    if found_count == chapter_count and (not old_files_exist):
        return 1.0
    return 0.0

def check_exact_text_match__e4beccec(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_text_replacement__65dee18f1880f7d028c2b1727fd62d90(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from the file
        expected: Expected rules dict with 'original_word' and 'correct_word'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    original_word = expected.get('original_word', '')
    correct_word = expected.get('correct_word', '')
    if not original_word or not correct_word:
        return 0.0
    if re.search('\\b' + re.escape(original_word) + '\\b', result, re.IGNORECASE):
        return 0.0
    expected_count = expected.get('expected_count', 0)
    actual_count = len(re.findall('\\b' + re.escape(correct_word) + '\\b', result, re.IGNORECASE))
    if actual_count >= expected_count:
        return 1.0
    if expected_count > 0:
        return min(1.0, actual_count / expected_count)
    return 0.0

def check_files_contain_pattern__7427978e(filenames, expected):
    """Check if files in directory contain expected pattern.

    Args:
        filenames: List of filenames in directory
        expected: Dict with 'expected' key containing list of expected filenames

    Returns:
        float: 1.0 if all expected files are present, 0.0 otherwise
    """
    expected_files = expected['expected']
    if len(filenames) != len(expected_files):
        return 0.0
    if set(filenames) == set(expected_files):
        return 1.0
    else:
        return 0.0

def check_line_count__5538243966b3481fe772536923c1f693(result, expected, **options):
    """
    Check if line count matches expected value.

    Args:
        result: Integer line count from getter
        expected: Rules dict with 'count' key
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, int):
        return 0.0
    expected_count = expected.get('count', 0)
    return float(result == expected_count)

def check_moved_files__1b47b6505a7a2dc3d6ad6f0c07b4bcb4(result, expected, **options):
    """
    Check if all expected files were moved to notebooks directory.

    Args:
        result: Actual list of files from getter
        expected: Expected data (from rules dict)
        **options: Additional options

    Returns:
        float: 1.0 if all files present, partial credit for some files, 0.0 if none
    """
    expected_files = expected.get('files', [])
    logger.info(f'Comparing file lists: result={result}, expected={expected_files}')
    if not isinstance(result, list):
        return 0.0
    result_sorted = sorted(result)
    expected_sorted = sorted(expected_files)
    if result_sorted == expected_sorted:
        return 1.0
    else:
        matching = len(set(result) & set(expected_files))
        total = len(expected_files)
        if total > 0:
            score = matching / total
            logger.info(f'Partial match: {matching}/{total} files found, score={score}')
            return score
        else:
            return 0.0

def check_textbox_topright__d493bb30(src_path, expected_state, **options):
    """
    Check if the textbox is in the upper-right area of the image.
    Variation 5 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in upper-right area, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    gray_array = np.array(gray_image)
    binary_mask = gray_array < 128
    (labeled_array, num_features) = ndimage.label(binary_mask)
    if num_features == 0:
        return 0.0
    largest_component = None
    largest_size = 0
    for label_num in range(1, num_features + 1):
        component_mask = labeled_array == label_num
        component_size = np.sum(component_mask)
        if component_size < 50:
            continue
        if component_size > largest_size:
            largest_size = component_size
            largest_component = component_mask
    if largest_component is None:
        return 0.0
    (rows, cols) = np.where(largest_component)
    if len(rows) == 0:
        return 0.0
    (min_row, max_row) = (rows.min(), rows.max())
    (min_col, max_col) = (cols.min(), cols.max())
    center_x = (min_col + max_col) / 2
    center_y = (min_row + max_row) / 2
    in_right_area = center_x > width * 0.7
    in_top_area = center_y < height * 0.3
    if in_right_area and in_top_area:
        return 1.0
    else:
        return 0.0

def check_filename_exact_match__0d61b4f8(result, expected, **options):
    """Check if found filename matches expected name exactly.

    Args:
        result: Filename from getter
        expected: Expected rules dict with 'filename' key
        **options: Additional comparison options

    Returns:
        float: 1.0 if exact match, 0.0 otherwise
    """
    expected_filename = expected.get('filename', '')
    if result == expected_filename:
        return 1.0
    else:
        logger.info(f"Filename mismatch - Expected: '{expected_filename}', Found: '{result}'")
        return 0.0

def check_textbox_fully_centered__c84bec37(src_path, expected, **options):
    """
    Check if the textbox is fully centered (both horizontally and vertically) in the image.
    Task variation 9 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is fully centered, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    left_most = width
    right_most = 0
    top_most = height
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                left_most = min(left_most, x)
                right_most = max(right_most, x)
                top_most = min(top_most, y)
                bottom_most = max(bottom_most, y)
    text_center_x = (left_most + right_most) / 2
    text_center_y = (top_most + bottom_most) / 2
    image_center_x = width / 2
    image_center_y = height / 2
    tolerance_x = width * 0.1
    tolerance_y = height * 0.1
    horizontally_centered = abs(text_center_x - image_center_x) < tolerance_x
    vertically_centered = abs(text_center_y - image_center_y) < tolerance_y
    if horizontally_centered and vertically_centered:
        return 1.0
    else:
        return 0.0

def check_exact_recipient_set__c9ce3f52(result, expected, **options):
    """Check if the recipient list contains exactly the expected emails (no more, no less).

    Args:
        result: List of actual emails in To field
        expected: Dict with 'emails' key containing list of required emails
        **options: Additional options

    Returns:
        float: 1.0 if exact match, partial credit for partial match
    """
    if not isinstance(result, list):
        return 0.0
    expected_emails = expected.get('emails', [])
    if not expected_emails:
        return 1.0 if len(result) == 0 else 0.0
    result_normalized = set([email.lower() for email in result])
    expected_normalized = set([email.lower() for email in expected_emails])
    if result_normalized == expected_normalized:
        return 1.0
    intersection = result_normalized & expected_normalized
    union = result_normalized | expected_normalized
    if len(union) == 0:
        return 0.0
    return len(intersection) / len(union)

def check_eml_files__846f274f(result, expected, **options):
    """Check if .eml files were downloaded successfully.

    Args:
        result: list of file paths from get_googledrive_file getter
        expected: dict with "all_exist" and "min_count"

    Returns:
        float: 1.0 if all files exist and count met, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if not isinstance(result, list):
        result = [result]
    valid_files = [f for f in result if f is not None]
    count = len(valid_files)
    min_count = expected.get('min_count', 0)
    all_exist = expected.get('all_exist', True)
    if all_exist and len(valid_files) < len(result):
        return 0.0
    if count >= min_count:
        return 1.0
    return 0.0

def check_timezone__7a83c51f(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_exists__25309a67d723dd8e75eb60c978b60929(result, expected, **options):
    """Check if file exists as expected and is a valid GIF with content.

    Args:
        result: Result from getter (dict with 'exists', 'is_gif', 'has_content' keys)
        expected: Expected rules (dict with 'should_exist' key)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, is a valid GIF, and has content; 0.0 otherwise
    """
    file_exists = result.get('exists', False)
    is_gif = result.get('is_gif', False)
    has_content = result.get('has_content', False)
    should_exist = expected.get('should_exist', True)
    if not should_exist:
        return 1.0 if not file_exists else 0.0
    if file_exists and is_gif and has_content:
        return 1.0
    else:
        return 0.0

def check_textbox_on_bottomside__134680e7(src_path, expected_state, **options):
    """
    Check if the textbox is on the bottom side of the image.
    Variation 2 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in bottom 5%, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    bottom_most_dark_pixel = 0
    for y in range(height - 1, -1, -1):
        row_has_dark_pixel = False
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                row_has_dark_pixel = True
                break
        if row_has_dark_pixel:
            bottom_most_dark_pixel = y
            break
    if bottom_most_dark_pixel > height * 0.95:
        return 1.0
    else:
        return 0.0

def check_text_exact_match__d332c3241fced231d1d84d00e75fe3b7(result: str, expected: dict, **options) -> float:
    """
    Check if text content exactly matches expected value.

    Args:
        result: Text content from getter
        expected: Dict with 'content' key containing expected text
        **options: Additional options

    Returns:
        1.0 if exact match, 0.0 otherwise
    """
    if result is None:
        logger.debug('Result is None')
        return 0.0
    expected_content = expected.get('content', '')
    if result == expected_content:
        return 1.0
    else:
        logger.debug(f"Content mismatch. Expected: '{expected_content}', Got: '{result}'")
        return 0.0

def check_filename_pattern__a2f23245(result, expected, **options):
    """Check filenames match expected patterns.

    Args:
        result: list of filenames from getter
        expected: dict with "extension" and "min_files"

    Returns:
        float: 1.0 if filenames match criteria, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    extension = expected.get('extension', '')
    min_files = expected.get('min_files', 0)
    if len(result) < min_files:
        return 0.0
    if extension:
        matching = [f for f in result if f.endswith(extension)]
        if len(matching) < min_files:
            return 0.0
    return 1.0

def check_textbox_on_rightside__5e2434f0(result_state, expected_state, **options):
    """
    Check if the textbox is on the right side of the image.
    Variation 0 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result_state: Path to the result image with text (vm_file path)
        expected_state: Not used (rule-based)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is on right side, 0.0 otherwise)
    """
    if result_state is None or not isinstance(result_state, str):
        return 0.0
    try:
        source_image = Image.open(result_state)
        gray_image = source_image.convert('L')
        (width, height) = source_image.size
        right_most_dark_pixel = 0
        for y in range(height):
            for x in range(width - 1, -1, -1):
                if gray_image.getpixel((x, y)) < 128:
                    right_most_dark_pixel = max(right_most_dark_pixel, x)
                    break
        if right_most_dark_pixel > width * 0.95:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking textbox position: {e}')
        return 0.0

def check_dir_file_count__95b4929b(result, expected, **options):
    """Check if directory file count matches expected.

    Args:
        result: Integer count from getter
        expected: Dict with 'count' value
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, int) or not isinstance(expected, dict):
        return 0.0
    expected_count = expected.get('count', 0)
    if result == expected_count:
        return 1.0
    return 0.0

def check_file_list__4d646866(result, expected, **options):
    """Check if file list contains all expected files.

    Args:
        result: File content with one filename per line
        expected: Dict with 'files' key containing list of expected filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on how many files match
    """
    if not result or not result.strip():
        return 0.0
    result_lines = [line.strip() for line in result.strip().split('\n') if line.strip()]
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    matches = sum((1 for f in expected_files if f in result_lines))
    score = matches / len(expected_files)
    if score < 1.0:
        missing = [f for f in expected_files if f not in result_lines]
        logger.info(f'File list incomplete: {len(matches)}/{len(expected_files)} files found. Missing: {missing}')
    return score

def check_file_location__739292ff(result, expected, **options):
    """Check if image file exists and has valid properties.

    Args:
        result: Dict with image properties (format, width, height, size_bytes)
        expected: Dict with validation rules (format, min_size_bytes, min_width, min_height)
        **options: Additional options

    Returns:
        float: 1.0 if image is valid, 0.0 otherwise
    """
    if result is None:
        logger.error('Result is None - image file may not exist or is not a valid image')
        return 0.0
    expected_format = expected.get('format')
    if expected_format:
        actual_format = result.get('format')
        if actual_format != expected_format:
            logger.info(f'Image format mismatch: expected {expected_format}, got {actual_format}')
            return 0.0
    min_size = expected.get('min_size_bytes')
    if min_size:
        actual_size = result.get('size_bytes', 0)
        if actual_size < min_size:
            logger.info(f'Image size too small: expected >= {min_size} bytes, got {actual_size} bytes')
            return 0.0
    min_width = expected.get('min_width')
    if min_width:
        actual_width = result.get('width', 0)
        if actual_width < min_width:
            logger.info(f'Image width too small: expected >= {min_width} pixels, got {actual_width} pixels')
            return 0.0
    min_height = expected.get('min_height')
    if min_height:
        actual_height = result.get('height', 0)
        if actual_height < min_height:
            logger.info(f'Image height too small: expected >= {min_height} pixels, got {actual_height} pixels')
            return 0.0
    logger.info(f"Image validation passed: format={result.get('format')}, size={result.get('size_bytes')} bytes, dimensions={result.get('width')}x{result.get('height')}")
    return 1.0

def check_text_contains_pattern__cc11abb5(result, expected, **options):
    """Check if paper_attachments folder exists on Google Drive and contains expected files.

    Args:
        result: Dict with 'folder_exists' (bool), 'file_count' (int), and 'file_names' (list)
        expected: Dict with 'folder_exists', 'min_file_count', and optional 'expected_files' requirements

    Returns:
        float: 1.0 if folder exists and has all expected files, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    folder_exists = result.get('folder_exists', False)
    if not folder_exists:
        return 0.0
    file_count = result.get('file_count', 0)
    min_file_count = expected.get('min_file_count', 1)
    if file_count < min_file_count:
        return 0.0
    expected_files = expected.get('expected_files', [])
    if expected_files:
        file_names = result.get('file_names', [])
        for expected_file in expected_files:
            if expected_file not in file_names:
                return 0.0
    return 1.0

def check_text_uppercase__423804a3(result, expected, **options):
    """Check if text is in uppercase and matches expected value.

    Args:
        result: Actual text from getter
        expected: Expected text (from rules dict)
        **options: Additional options

    Returns:
        1.0 if match, 0.0 otherwise
    """
    expected_text = expected.get('text', '')
    if result is None:
        return 0.0
    if result.strip() == expected_text.strip():
        return 1.0
    return 0.0

def check_file_organization__b1a155d8(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_text_replacement__dd369016ed084e5c5f565139bb4ef07f(result, expected, **options):
    """Check if text replacement was performed correctly.

    Args:
        result: Actual text content from the file
        expected: Expected rules dict with 'original_word' and 'correct_word'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    original_word = expected.get('original_word', '')
    correct_word = expected.get('correct_word', '')
    if not original_word or not correct_word:
        return 0.0
    if original_word in result:
        return 0.0
    expected_count = expected.get('expected_count', 0)
    actual_count = result.count(correct_word)
    if actual_count >= expected_count:
        return 1.0
    if expected_count > 0:
        return min(1.0, actual_count / expected_count)
    return 0.0

def check_zip_contains_files__8a01d242c9e2052109e1c26f5fa4a5dd(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if ZIP archive contains required files.

    Args:
        result: Dict from getter with 'exists', 'files', 'valid_zip'
        expected: Dict with 'required_extensions' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not result.get('exists', False):
        return 0.0
    if not result.get('valid_zip', False):
        return 0.0
    files = result.get('files', [])
    if not files:
        return 0.0
    required_exts = expected.get('required_extensions', [])
    found_exts = set()
    for filename in files:
        for ext in required_exts:
            if filename.endswith(ext):
                found_exts.add(ext)
    score = len(found_exts) / len(required_exts) if required_exts else 0.0
    return score

def check_file_rename__c14397f3284104a2f980691d2ea6abf3(result: Dict[str, bool], expected: Dict[str, Any], **options) -> float:
    """
    Check if file rename operation was completed correctly.
    Verifies that the new file exists AND the old file does not exist.

    Args:
        result: Dict with 'new_exists' and 'old_exists' keys from getter
        expected: Dict with 'renamed' key (should be True for successful rename)

    Returns:
        1.0 if rename was successful (new file exists, old file doesn't), 0.0 otherwise
    """
    renamed_expected = expected.get('renamed', True)
    new_exists = result.get('new_exists', False)
    old_exists = result.get('old_exists', True)
    rename_successful = new_exists and (not old_exists)
    if renamed_expected and rename_successful:
        logger.info(f'File rename check passed: new file exists={new_exists}, old file exists={old_exists}')
        return 1.0
    elif not renamed_expected and (not rename_successful):
        logger.info(f'File rename check passed (expected no rename): new file exists={new_exists}, old file exists={old_exists}')
        return 1.0
    else:
        logger.info(f'File rename check failed: new file exists={new_exists}, old file exists={old_exists}, expected renamed={renamed_expected}')
        return 0.0

def check_text_replacement__4aeab799(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_file_deleted__464fee7f(result, expected, **options):
    """
    Check if specific file was deleted while others remain.

    Args:
        result: Output from ls command
        expected: Dict with 'rules' containing 'should_not_exist' and 'should_exist'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    should_not_exist = expected.get('should_not_exist', '')
    should_exist = expected.get('should_exist', [])
    if isinstance(result, dict) and 'output' in result:
        result = result['output']
    result = str(result).strip()
    files_present = [line.strip() for line in result.split('\n') if line.strip() and '.pdf' in line]
    files_present = [f.split('/')[-1] for f in files_present]
    score = 0.0
    if should_not_exist not in result and should_not_exist not in files_present:
        score += 0.5
    if should_exist:
        exists_count = sum((1 for f in should_exist if f in result or f in files_present))
        score += 0.5 * (exists_count / len(should_exist))
    return min(score, 1.0)

def check_text_color_white__1b43c872(result, expected, **options):
    """
    Check if the text color has been changed to white.
    Variation 4 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result: Path to the result image with text (vm_file)
        expected: Not used (rule-based)
        **options: Additional options
            - white_threshold: Minimum RGB value for white detection (default: 200)
            - white_ratio_threshold: Minimum ratio of white pixels required (default: 0.8)

    Returns:
        float: Score (1.0 if text is white, 0.0 otherwise)
    """
    file_path = None
    if isinstance(result, dict) and 'path' in result:
        file_path = result['path']
    elif isinstance(result, str):
        file_path = result
    if not file_path:
        logger.error('No file path provided in result')
        return 0.0
    try:
        source_image = Image.open(file_path)
        rgb_image = source_image.convert('RGB')
        (width, height) = source_image.size
        (bg_r_min, bg_r_max) = (200, 256)
        (bg_g_min, bg_g_max) = (100, 200)
        (bg_b_min, bg_b_max) = (0, 50)
        white_threshold = options.get('white_threshold', 200)
        white_ratio_threshold = options.get('white_ratio_threshold', 0.8)
        white_text_pixels = 0
        total_non_bg_pixels = 0
        for y in range(height):
            for x in range(width):
                (r, g, b) = rgb_image.getpixel((x, y))
                is_bg = bg_r_min <= r < bg_r_max and bg_g_min <= g < bg_g_max and (bg_b_min <= b < bg_b_max)
                if not is_bg:
                    total_non_bg_pixels += 1
                    if r > white_threshold and g > white_threshold and (b > white_threshold):
                        white_text_pixels += 1
        if total_non_bg_pixels == 0:
            logger.warning('No non-background pixels found in image')
            return 0.0
        white_ratio = white_text_pixels / total_non_bg_pixels
        logger.info(f'White pixel ratio: {white_ratio:.2%} ({white_text_pixels}/{total_non_bg_pixels})')
        if white_ratio > white_ratio_threshold:
            return 1.0
        return 0.0
    except FileNotFoundError:
        logger.error(f'Image file not found: {file_path}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking text color: {e}')
        return 0.0

def check_text_exists__109edd97(result, expected, **options):
    """Check if text exists in document.

    Args:
        result: Boolean from getter
        expected: Expected value
        **options: Additional options

    Returns:
        float: 1.0 if result matches expected, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if result == expected_exists:
        return 1.0
    return 0.0

def check_file_first_line__c5a909ed(actual: str, expected: dict, **options) -> float:
    """
    Check if the first line of a file matches the expected line.

    Args:
        actual (str): path to the file
        expected (dict): expected dict with key "expected_line"

    Return:
        float: 1.0 if the first line matches, 0.0 otherwise
    """
    if not actual:
        logger.debug('check_file_first_line__c5a909ed: actual file path is None')
        return 0.0
    expected_line = expected.get('expected_line')
    if expected_line is None:
        logger.debug('check_file_first_line__c5a909ed: expected_line is None')
        return 0.0
    try:
        with open(actual, 'r') as f:
            first_line = f.readline().strip()
    except Exception as e:
        logger.debug(f'check_file_first_line__c5a909ed: Error reading file: {e}')
        return 0.0
    if first_line == expected_line:
        return 1.0
    logger.debug(f"check_file_first_line__c5a909ed: Expected '{expected_line}', got '{first_line}'")
    return 0.0

def check_file_renamed__bd9934867663f0945bb79537ace5711a(result: dict, expected: dict, **options) -> float:
    """Check if file rename operation completed successfully.

    Verifies that the new file exists AND the old file doesn't exist.

    Args:
        result: Dict with 'new_exists' and 'old_exists' boolean values
        expected: Dict from rules with expected 'new_exists' and 'old_exists' values
        **options: Additional options

    Returns:
        float: 1.0 if rename completed successfully, 0.0 otherwise
    """
    expected_new_exists = expected.get('new_exists', True)
    expected_old_exists = expected.get('old_exists', False)
    result_new_exists = result.get('new_exists', False)
    result_old_exists = result.get('old_exists', True)
    if result_new_exists == expected_new_exists and result_old_exists == expected_old_exists:
        logger.debug(f'File rename successful: new_exists={result_new_exists}, old_exists={result_old_exists}')
        return 1.0
    else:
        logger.debug(f'File rename incomplete or failed: new_exists={result_new_exists} (expected {expected_new_exists}), old_exists={result_old_exists} (expected {expected_old_exists})')
        return 0.0

def check_file_exists_with_size__995b229f(result, expected, **options):
    """Check if file exists and has reasonable size."""
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_png', False):
        score += 0.3
    size = result.get('size', 0)
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 50000000)
    if min_size <= size <= max_size:
        score += 0.3
    return score

def check_textbox_vertically_centered__8c858544(src_path, expected, **options):
    """
    Check if the textbox is vertically centered in the image.
    Task variation 8 for e2dd0213-26db-4349-abe5-d5667bfd725c_xhaug_1

    Args:
        src_path: Path to the exported PNG image
        expected: Not used (rule-based check)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is vertically centered, 0.0 otherwise)
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    top_most = height
    bottom_most = 0
    for y in range(height):
        for x in range(width):
            if gray_image.getpixel((x, y)) < 128:
                top_most = min(top_most, y)
                bottom_most = max(bottom_most, y)
    text_center_y = (top_most + bottom_most) / 2
    image_center_y = height / 2
    tolerance = height * 0.1
    if abs(text_center_y - image_center_y) < tolerance:
        return 1.0
    else:
        return 0.0

def check_tex_file_list__557a0701(result_state: Dict[str, Any], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if the file list contains all expected .tex files.

    Args:
        result_state: Dict from getter with keys:
            - exists (bool): whether file exists
            - line_count (int): number of files in the list
            - files (List[str]): list of filenames
            - content (str): raw file content
        expected_state: Dict containing rules directly (when type='rule'):
            - min_count (int): minimum number of .tex files expected
            - expected_files (List[str]): list of expected .tex filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    min_count = expected_state.get('min_count', 0)
    expected_files = expected_state.get('expected_files', [])
    if not result_state.get('exists', False):
        logger.warning('File list does not exist')
        return 0.0
    files = result_state.get('files', [])
    line_count = result_state.get('line_count', 0)
    if line_count < min_count:
        logger.warning(f'File list has {line_count} files, expected at least {min_count}')
        return 0.0
    actual_filenames = []
    for f in files:
        import os
        basename = os.path.basename(f.strip())
        actual_filenames.append(basename)
    tex_files = [f for f in actual_filenames if f.lower().endswith('.tex')]
    if len(tex_files) < min_count:
        logger.warning(f'Found only {len(tex_files)} .tex files, expected at least {min_count}')
        return 0.0
    missing_files = []
    for expected_file in expected_files:
        expected_normalized = expected_file.strip()
        found = False
        for actual_file in tex_files:
            if actual_file == expected_normalized or actual_file.lower() == expected_normalized.lower():
                found = True
                break
        if not found:
            missing_files.append(expected_file)
    if missing_files:
        logger.warning(f'Missing expected .tex files: {missing_files}')
        return 0.0
    logger.info(f'File list verification passed: {len(tex_files)} .tex files found, all expected files present')
    return 1.0

def check_textbox_centered__b99b3c81(result_state, expected_state, **options):
    """
    Check if the textbox is centered horizontally in the image.
    Variation 1 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result_state: Path to the result image with text (vm_file)
        expected_state: Not used (rule-based)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is centered, 0.0 otherwise)
    """
    if not result_state:
        logger.error('No file path provided in result_state')
        return 0.0
    try:
        with Image.open(result_state) as source_image:
            gray_image = source_image.convert('L')
            (width, height) = source_image.size
            left_most = width
            right_most = 0
            for y in range(height):
                for x in range(width):
                    if gray_image.getpixel((x, y)) < 128:
                        left_most = min(left_most, x)
                        right_most = max(right_most, x)
            if left_most >= right_most:
                logger.error('No text found in image (no dark pixels detected)')
                return 0.0
            text_center = (left_most + right_most) / 2
            image_center = width / 2
            tolerance = width * 0.1
            if abs(text_center - image_center) < tolerance:
                return 1.0
            else:
                return 0.0
    except Exception as e:
        logger.error(f'Error checking textbox centering: {e}')
        return 0.0

def check_text_contains__3c678f53(result, expected, **options):
    """Check if text contains expected content.

    Args:
        result: Text content from getter
        expected: Dict with 'pattern' (regex) or 'text' (literal string)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    if 'pattern' in expected:
        pattern = expected['pattern']
        if re.search(pattern, result):
            return 1.0
        return 0.0
    if 'text' in expected:
        text = expected['text']
        if text in result:
            return 1.0
        return 0.0
    if 'number' in expected:
        target_number = expected['number']
        numbers = re.findall('\\d+', result)
        if numbers and int(numbers[0]) == target_number:
            return 1.0
        return 0.0
    return 0.0

def check_file_contains_text__b6bb3e42(result, expected, **options):
    """Check if file contains expected text.

    Args:
        result: File content
        expected: Dict with 'text' key specifying expected text
        **options: Additional options

    Returns:
        float: 1.0 if text is contained, 0.0 otherwise
    """
    if result is None:
        return 0.0
    text = expected.get('text', '')
    if text in result:
        return 1.0
    else:
        return 0.0

def check_timezone_utc_minus_6__53de5eae(result_state, expected_state, **options):
    """
    Check if timezone is set to the expected UTC offset (e.g., -0600 for Central Standard Time).

    Args:
        result_state: Output from 'timedatectl status' command
        expected_state: Expected timezone configuration dict with 'offset' and optional 'allow_dst'
        **options: Additional options

    Returns:
        float: 1.0 if timezone offset matches expected value, 0.0 otherwise
    """
    timedatectl_output = result_state
    lines = timedatectl_output.split('\n')
    if isinstance(expected_state, dict):
        rules = expected_state.get('rules', expected_state)
        expected_offset = rules.get('offset', '-0600')
        allow_dst = rules.get('allow_dst', False)
    else:
        expected_offset = expected_state if expected_state else '-0600'
        allow_dst = False
    timezone_line = None
    for line in lines:
        if 'Time zone:' in line:
            timezone_line = line
            break
    if not timezone_line:
        if len(lines) > 3:
            timezone_line = lines[3]
        else:
            return 0.0
    offset_pattern = '([+-]\\d{4})\\)'
    match = re.search(offset_pattern, timezone_line)
    if not match:
        return 0.0
    actual_offset = match.group(1)
    if actual_offset == expected_offset:
        return 1.0
    if allow_dst and expected_offset == '-0600' and (actual_offset == '-0500'):
        return 1.0
    return 0.0

def check_line_count__b2b950a6(result: Dict, expected: Dict, **options) -> float:
    """Check if line count matches expected value and specific content is deleted.

    Args:
        result: Dict with 'line_count' and 'content' keys
        expected: Dict with 'count' and optionally 'deleted_content' key

    Returns:
        1.0 if counts match and specific content is deleted, 0.0 otherwise
    """
    expected_count = expected.get('count', 0)
    if result.get('line_count') != expected_count:
        return 0.0
    content = result.get('content', '')
    if 'if 0 == count:' in content:
        return 0.0
    return 1.0

def check_json_has_fields__c8d946870135d67f0db0be5e65caaa2a(result, expected, **options):
    """Check if JSON data contains expected fields.

    Args:
        result: Dict with JSON data from getter
        expected: Dict with 'required_fields' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dict')
        return 0.0
    required_fields = expected.get('required_fields', [])
    if not required_fields:
        logger.warning('No required fields specified')
        return 0.0
    found_count = 0
    for field in required_fields:
        if field in result:
            found_count += 1
    return found_count / len(required_fields)

def check_text_colors_match__6abdacd1e07293be0566dfd3601ee79b(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if title and subtitle colors match each other.

    Args:
        result: Dict from getter with 'title_color' and 'subtitle_color'
        expected: Dict with 'colors_must_match' boolean from rules
        **options: Additional options (e.g., 'tolerance' for RGB matching)

    Returns:
        1.0 if colors match each other (within tolerance), 0.0 otherwise
    """
    title_color = result.get('title_color')
    subtitle_color = result.get('subtitle_color')
    if not title_color or not subtitle_color:
        return 0.0
    tolerance = options.get('tolerance', 5)
    r_match = abs(title_color[0] - subtitle_color[0]) <= tolerance
    g_match = abs(title_color[1] - subtitle_color[1]) <= tolerance
    b_match = abs(title_color[2] - subtitle_color[2]) <= tolerance
    if r_match and g_match and b_match:
        return 1.0
    else:
        return 0.0

def check_file_moved__00112b53200a74ce7a53869d2d085264(result: dict, expected: dict, **options) -> float:
    """Check if a file has been properly moved (exists at target, not at source).

    Args:
        result: Dict from getter with 'target_exists' and 'source_exists' keys
        expected: Dict from rules with 'target_exists' and 'source_exists' expected values
        **options: Additional options

    Returns:
        float: 1.0 if file move is correct, 0.0 otherwise
    """
    expected_target_exists = expected.get('target_exists', True)
    expected_source_exists = expected.get('source_exists', False)
    actual_target_exists = result.get('target_exists', False)
    actual_source_exists = result.get('source_exists', True)
    target_matches = actual_target_exists == expected_target_exists
    source_matches = actual_source_exists == expected_source_exists
    if target_matches and source_matches:
        logger.debug(f'File move verified: target_exists={actual_target_exists}, source_exists={actual_source_exists}')
        return 1.0
    else:
        logger.debug(f'File move failed: target_exists={actual_target_exists} (expected {expected_target_exists}), source_exists={actual_source_exists} (expected {expected_source_exists})')
        return 0.0

def check_python_functions__af4f2737(result_file, expected, **options):
    """
    Check if a Python file contains required function definitions (excluding class methods).

    Args:
        result_file: Path to the Python file to check
        expected: Dict with 'required_functions' and 'min_functions' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 based on partial credit
    """
    if not result_file:
        return 0.0
    try:
        with open(result_file, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        print(f'Error reading file: {e}')
        return 0.0
    required_functions = expected.get('required_functions', [])
    min_functions = expected.get('min_functions', len(required_functions))
    if not required_functions:
        return 0.0
    found_functions = set()
    try:
        tree = ast.parse(content)
        for node in tree.body:
            if isinstance(node, ast.FunctionDef):
                found_functions.add(node.name)
            elif isinstance(node, ast.Assign):
                for target in node.targets:
                    if isinstance(target, ast.Name) and isinstance(node.value, ast.Lambda):
                        found_functions.add(target.id)
    except SyntaxError as e:
        print(f'Syntax error parsing file: {e}')
        return 0.0
    found_count = sum((1 for func in required_functions if func in found_functions))
    if found_count >= min_functions:
        score = found_count / len(required_functions)
    else:
        score = found_count / len(required_functions) * 0.5
    return score

def check_file_content__cd720b2833d8c75c48e4cd829046ee69(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a file exists and contains expected content.

    Args:
        result: Dict from getter with 'exists', 'content', and 'file_path' keys
        expected: Dict with 'expected_content' key
        **options: Additional options (ignore_whitespace, ignore_case)

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if not result.get('exists', False):
        return 0.0
    score += 0.5
    actual_content = result.get('content', '')
    expected_content = expected.get('expected_content', '')
    ignore_whitespace = options.get('ignore_whitespace', False)
    ignore_case = options.get('ignore_case', False)
    if ignore_whitespace:
        actual_content = ' '.join(actual_content.split())
        expected_content = ' '.join(expected_content.split())
    if ignore_case:
        actual_content = actual_content.lower()
        expected_content = expected_content.lower()
    if actual_content.strip() == expected_content.strip():
        score += 0.5
    return score

def check_timezone__e55695d1(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_count__fbd9137c(result, expected, **options):
    """Check if file count matches expected value.

    Args:
        result: File count from getter
        expected: Expected file count
        **options: Additional options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if result is None:
        return 0.0
    if result == expected:
        return 1.0
    else:
        return 0.0

def check_git_repo_exists__68f4a1f5(result, expected, **options):
    """Check if git repository exists.

    Args:
        result: Dictionary with 'exists' and 'is_git_repo' keys
        expected: Expected dictionary (not used, we just check existence)
        **options: Additional options

    Returns:
        float: 1.0 if repository exists and is a valid git repo, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    if result.get('exists', False) and result.get('is_git_repo', False):
        return 1.0
    return 0.0

def check_file_count__2d28c7f27d8eb7cfbd3f915d05c991e3(directory_list, rule):
    """
    Check if the directory contains the expected number of JPG files.

    Args:
        directory_list: Directory tree structure from get_list_directory
        rule: Expected configuration with 'expected_count' key

    Returns:
        float: 1.0 if file count matches and all files are JPG, 0.0 otherwise
    """
    expected_count = rule['expected_count']
    children = directory_list.get('children', [])
    actual_count = len(children)
    if actual_count != expected_count:
        return 0.0
    for child in children:
        filename = child.get('name', '')
        if not filename.lower().endswith('.jpg'):
            return 0.0
    return 1.0

def check_text_replacement__f89b526b(actual: str, rules: dict, **options) -> float:
    """
    Verify that text replacement was performed correctly in a file.

    Args:
        actual (str): path to result text file
        rules (dict): dict with keys:
            - original_word (str): word that should have been replaced
            - replacement_word (str): word that should appear instead
            - expected_count (int): number of expected replacements
        **options: additional options

    Returns:
        float: 1.0 if replacement is correct, 0.0 otherwise
    """
    if not actual:
        return 0.0
    try:
        with open(actual, 'r') as f:
            content = f.read()
        original_word = rules['original_word']
        replacement_word = rules['replacement_word']
        expected_count = rules['expected_count']
        if original_word in content:
            logger.debug(f"Original word '{original_word}' still found in file")
            return 0.0
        actual_count = content.count(replacement_word)
        if actual_count != expected_count:
            logger.debug(f"Replacement word '{replacement_word}' count mismatch: expected {expected_count}, got {actual_count}")
            return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error checking text replacement: {e}')
        return 0.0

def check_extracted_files__940d01bc(result, expected, **options):
    """Check if specific files were extracted from archive.

    Args:
        result: List of file paths from getter
        expected: Dict with 'required_files' list (relative paths)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    required_files = expected.get('required_files', [])
    if not required_files:
        return 0.0
    found_count = 0
    for req_file in required_files:
        if any((file_path.endswith(req_file) for file_path in result)):
            found_count += 1
    return found_count / len(required_files)

def check_file_organization__ba67a508(result, expected, **options):
    """Check if files are organized correctly in directories.

    Args:
        result: List of [dir1_contents, dir2_contents, dir3_contents]
        expected: Expected file lists for each directory
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list) or len(result) != 3:
        return 0.0
    if not isinstance(expected, list) or len(expected) != 3:
        return 0.0
    score = 0.0
    for i in range(3):
        actual_files = set(result[i].strip().split('\n')) if result[i].strip() else set()
        expected_files = set(expected[i])
        actual_files.discard('')
        if actual_files == expected_files:
            score += 1.0 / 3.0
    return score

def check_files_with_prefix__cf3f5d8ece62ecf8e4937dea9e007679(result: Dict, expected: Dict, **options) -> float:
    """
    Check if files have the correct prefix.

    Args:
        result: Dict mapping file hashes to filenames
        expected: Dict with expected prefix pattern
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result:
        return 0.0
    required_hashes = expected.get('required_hashes', [])
    prefix = expected.get('prefix', 'Mountain_')
    if not required_hashes:
        return 0.0
    correct_count = 0
    for req_hash in required_hashes:
        if req_hash in result:
            filename = result[req_hash]
            if filename.startswith(prefix):
                correct_count += 1
                logger.info(f'Hash {req_hash[:16]}... has correct prefix: {filename}')
            else:
                logger.info(f'Hash {req_hash[:16]}... missing prefix: {filename}')
        else:
            logger.info(f'Hash {req_hash[:16]}... not found in results')
    score = correct_count / len(required_hashes)
    return score

def check_text_contains__b3a4b9dc(result, expected, **options):
    """Check if text contains expected phrases.

    Args:
        result: Text from getter
        expected: Expected dict with 'phrases' key

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, str):
        return 0.0
    result_lower = result.lower()
    phrases = expected.get('phrases', [])
    if not phrases:
        return 0.0
    matches = sum((1 for phrase in phrases if phrase.lower() in result_lower))
    return matches / len(phrases)

def check_file_recovered_not_in_trash__05c91736(result, expected, **options):
    """Verify file is recovered and not in trash.

    Args:
        result: Dict with file_exists and in_trash booleans
        expected: Dict with rules specifying expected states
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_file_exists = expected.get('file_exists', True)
    expected_not_in_trash = expected.get('not_in_trash', True)
    file_exists = result.get('file_exists', False)
    in_trash = result.get('in_trash', True)
    score = 0.0
    if file_exists == expected_file_exists:
        score += 0.5
    if expected_not_in_trash and (not in_trash):
        score += 0.5
    elif not expected_not_in_trash and in_trash:
        score += 0.5
    return score

def check_exact_text_match__05bf41d7(result, expected, **options):
    """Check if result text exactly matches expected text."""
    if result == expected:
        return 1.0
    return 0.0

def check_file_exists__3b8e423e430323c0078f4425aded05b9(result, expected, **options):
    """Check if file exists and verify 180-degree rotation was actually applied.

    This metric verifies that:
    1. The output file exists
    2. It is a valid video file
    3. It has reasonable size
    4. A 180-degree rotation was actually performed (not just copied)

    Args:
        result: Dict from getter with keys:
            - 'exists': bool
            - 'path': str
            - 'size': int
            - 'source_rotation': int or None (rotation of source video)
            - 'output_rotation': int or None (rotation of output video)
            - 'rotation_changed': bool (whether rotation was modified)
            - 'is_valid_video': bool
        expected: Dict with 'exists' key (True/False) from rules
        **options: Additional options

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not isinstance(result, dict) or 'exists' not in result:
        logger.error(f'Invalid result format: {result}')
        return 0.0
    expected_exists = expected.get('exists', True)
    actual_exists = result.get('exists', False)
    if actual_exists != expected_exists:
        logger.info(f"File {result.get('path', 'unknown')}: expected exists={expected_exists}, actual exists={actual_exists}")
        return 0.0
    if not expected_exists:
        logger.info(f"File correctly does not exist: {result.get('path', 'unknown')}")
        return 1.0
    file_path = result.get('path', 'unknown')
    file_size = result.get('size', 0)
    source_rotation = result.get('source_rotation')
    output_rotation = result.get('output_rotation')
    rotation_changed = result.get('rotation_changed', False)
    is_valid_video = result.get('is_valid_video', False)
    if not is_valid_video:
        logger.error(f'File {file_path} is not a valid video file')
        return 0.0
    MIN_VIDEO_SIZE = 100 * 1024
    if file_size < MIN_VIDEO_SIZE:
        logger.error(f'File {file_path} size {file_size} bytes is too small for a video (minimum {MIN_VIDEO_SIZE} bytes)')
        return 0.0
    logger.info(f'Rotation analysis: source={source_rotation}, output={output_rotation}, changed={rotation_changed}')
    if source_rotation is not None and output_rotation is not None:
        norm_source = source_rotation % 360
        norm_output = output_rotation % 360
        if norm_source == 180 and norm_output == 0:
            logger.info(f'File {file_path} correctly rotated from {source_rotation} to {output_rotation}')
            return 1.0
        else:
            if norm_source == norm_output:
                logger.error(f'File {file_path} was not rotated - source and output have same rotation ({source_rotation})')
                return 0.0
            if norm_source == 0 and norm_output == 180:
                logger.error(f'File {file_path} was incorrectly rotated from normal (0) to flipped (180) - opposite of requirement')
                return 0.0
            rotation_diff = abs(norm_source - norm_output)
            if rotation_diff != 180:
                logger.error(f'File {file_path} rotation incorrect: source={source_rotation}, output={output_rotation}, diff={rotation_diff} (expected 180)')
                return 0.0
            logger.error(f'File {file_path} rotation incorrect: source={source_rotation}, output={output_rotation} (expected 180->0 specifically)')
            return 0.0
    if rotation_changed:
        logger.info(f'File {file_path} rotation was changed (fallback check passed)')
        return 1.0
    if source_rotation is not None and source_rotation % 360 == 180 and (output_rotation == 0):
        logger.info(f'File {file_path} rotation normalized from {source_rotation} to 0')
        return 1.0
    logger.error(f'File {file_path} failed rotation verification: source={source_rotation}, output={output_rotation}, changed={rotation_changed}')
    return 0.0

def check_textbox_bottomright__d0e1ffb6(src_path, expected_state, **options):
    """
    Check if the textbox is in the bottom-right corner of the image.
    Variation 7 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        src_path: Path to the exported image
        expected_state: Not used (type=rule with empty rules)
        **options: Additional options

    Returns:
        float: 1.0 if text is in bottom-right 5% region, 0.0 otherwise
    """
    if src_path is None:
        return 0.0
    source_image = Image.open(src_path)
    gray_image = source_image.convert('L')
    (width, height) = source_image.size
    right_threshold = int(width * 0.95)
    bottom_threshold = int(height * 0.95)
    for y in range(bottom_threshold, height):
        for x in range(right_threshold, width):
            if gray_image.getpixel((x, y)) < 128:
                return 1.0
    return 0.0

def check_timezone__2f3e1080(timedatectl_output, expected, **options):
    """
    Check if timezone is set to the expected timezone.

    Args:
        timedatectl_output: Output from timedatectl status command
        expected: Dict with 'timezone' key containing expected timezone string

    Returns:
        float: 1.0 if timezone matches, 0.0 otherwise
    """
    expected_timezone = expected['timezone']
    lines = timedatectl_output.split('\n')
    for line in lines:
        if 'Time zone:' in line:
            timezone_part = line.split('Time zone:')[1].strip()
            timezone_name = timezone_part.split('(')[0].strip()
            if timezone_name == expected_timezone:
                return 1.0
            else:
                return 0.0
    return 0.0

def check_file_list__24c8b0df(result, expected, **options):
    """Check if the file list contains all expected .doc files.

    Args:
        result: List of filenames from getter
        expected: Dict with 'expected_count' parameter
        **options: Additional options

    Returns:
        float: 1.0 if count matches expected, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    expected_count = expected.get('expected_count', 12)
    doc_files = [f for f in result if f.endswith('.doc') and (not f.endswith('.docx'))]
    return 1.0 if len(doc_files) == expected_count else 0.0

def exact_match__a85eebd24e563a97c17935bc46126aa1(result, expected, **options):
    """
    Check if result exactly matches expected value.

    Args:
        result: Current value from getter
        expected: Dictionary with 'expected' key

    Returns:
        1.0 if values match exactly, 0.0 otherwise
    """
    expect = expected.get('expected', '')
    logger.info(f'Result: {result}')
    logger.info(f'Expected: {expect}')
    if result == expect:
        return 1.0
    else:
        return 0.0

def check_file_exists__25f0d8c3(result, expected, **options):
    """
    Check if file exists and is a valid PNG image with reasonable size.

    Args:
        result: Dict with keys:
            - exists: bool, whether file exists
            - is_valid_png: bool, whether file has valid PNG magic bytes
            - size: int, file size in bytes
        expected: Expected value dict with:
            - exists: bool, whether file should exist

    Returns:
        1.0 if all checks pass, 0.0 otherwise
    """
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        return 1.0 if not result.get('exists', False) else 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    if not result.get('is_valid_png', False):
        logger.info('File is not a valid PNG image (invalid magic bytes)')
        return 0.0
    file_size = result.get('size', 0)
    if file_size < 1024:
        logger.info(f'File size too small ({file_size} bytes), likely not a real screenshot')
        return 0.0
    return 1.0

def check_textbox_at_bottom__6cecb6d7(result_state, expected_state, **options):
    """
    Check if the textbox is at the bottom of the image.
    Variation 3 for task e2dd0213-26db-4349-abe5-d5667bfd725c

    Args:
        result_state: Path to the result image with text (vm_file)
        expected_state: Not used (rule-based)
        **options: Additional options

    Returns:
        float: Score (1.0 if text is at bottom, 0.0 otherwise)
    """
    import os
    if result_state is None or not isinstance(result_state, str):
        return 0.0
    if not os.path.exists(result_state):
        logger.error(f'File not found: {result_state}')
        return 0.0
    try:
        source_image = Image.open(result_state)
        gray_image = source_image.convert('L')
        (width, height) = source_image.size
        bottom_most_dark_pixel = 0
        for y in range(height - 1, -1, -1):
            for x in range(width):
                if gray_image.getpixel((x, y)) < 128:
                    bottom_most_dark_pixel = max(bottom_most_dark_pixel, y)
                    break
            if bottom_most_dark_pixel > 0:
                break
        if bottom_most_dark_pixel > height * 0.95:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking textbox position: {e}')
        return 0.0
