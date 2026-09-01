"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['check_qt_fullscreen_control__f9958b0e', 'check_vlc_fullscreen_state__d7a7afc9', 'check_vlc_paused__60670fd49b57e833635bf052410c24fc', 'check_vlc_playback_state__fc6f4242', 'check_vlc_loop__9eea6b3721652abbf42884414d2268c7', 'check_vlc_repeat_mode__9d232fb9', 'check_vlc_qt_startminimized__5a4867ff', 'check_vlc_subtitle_autodetect__d8ae66c5', 'check_vlc_skip_frames__99a2f14e8a82b525ea717f160309e108', 'check_vlc_screenshot_format__5940c66b', 'check_qt_autosave_volume__9190a44d', 'check_vlc_loop_mode__23c7a390', 'check_vlc_video_effects__84f74170', 'check_vlc_playlist_tree__2b678efeda2f5b06f01e1ca5035b98c7', 'check_vlc_always_on_top__850d1898', 'check_vlc_snapshot_format__772d6c83', 'check_vlc_bookmark__1c6b323b', 'check_vlc_stopped__c96690ce', 'check_vlc_recent_media__aa13955b1d571b5f9f91c4b58d07755a', 'check_vlc_video_transformed__78c019ea', 'check_qt_playlist_tree__bf5eef29', 'check_vlc_snapshot_exists__57da68f2', 'check_mp3_file_valid__ce2e74ee9437f284464037da6df4e453', 'check_vlc_subtitle_folder__73747be5', 'check_vlc_screenshot_path__fe752f70', 'check_vlc_paused__8f5f9872b32c7843d08f183a299b915a', 'check_vlc_screenshot_exists__167faa79', 'check_ogg_audio_extraction__c03bf6861f4e39f709595e582520ac0f', 'check_qt_autosave_volume__f8ad1f01', 'check_qt_video_autoresize__a74697ad', 'check_video_on_top__e48d3a20', 'check_mp3_count__be089f93', 'check_vlc_repeat_loop_enabled', 'check_vlc_video_transformed__a899663b', 'check_qt_video_autoresize__b7688244', 'check_vlc_audio_output__c4a2eaec', 'check_vlc_recordings_folder__055ce0bbe52cb194135b41b067eca986', 'check_qt_updates_notif__f6b51378', 'check_vlc_volume__9cc567291e2f7ec7556a7b0f92f9b42b', 'check_snapshot_format__ccaf0dcd', 'check_vlc_bookmarks_created__18214e12ab9f5ae2ef34a7b82632e6f9', 'check_vlc_playback_rate__c979999ccfa7b4560c7eb8656e505ff9', 'check_qt_privacy_ask__8eb13985', 'check_audio_extracted__dde7bdf2', 'check_qt_privacy_ask__132d5c78', 'is_vlc_playing_youtube__bba3381f', 'check_vlc_video_transformed__cb59ca2e', 'check_vlc_paused__df488c66748635e65016c0e9e7cc92c2', 'check_vlc_screenshot_folder__3d380bd3', 'check_vlc_qt_pause_minimized__7c31b71a', 'check_vlc_minimized__667c58f7', 'check_vlc_muted__cf575cad41cf5daac3aae885ce8cfe2a', 'check_vlc_snapshot_format__5ed539c67601c57bd107b5d8b04526be', 'check_vlc_video_on_top__9ca0ff77c9570a9a6e2f317bec352471', 'check_vlc_aout__a69f9ac550f1e14ddbeccea862a27f73', 'check_vlc_loop_mode__78fd4227', 'check_qt_continue__a10b6972', 'check_vlc_aspect_ratio__d1470170', 'check_audio_file_created__4071abfac3d4e1d773938a9f2b9279b5', 'check_vlc_playback_speed__c3b737f4', 'check_vlc_video_transformed__636928c8', 'check_vlc_subtitle_autoload__989833e1', 'check_vlc_subtitle_setting__71783eb7', 'check_vlc_bookmarks_count__3288b82e', 'check_vlc_video_transformed__2331e840', 'check_vlc_subtitle_disabled__65753b1a', 'check_vlc_repeat_mode__8e82e8fa6a0ee6159e766778349645ac', 'check_qt_start_minimized__b23ac353', 'check_vlc_video_output_dir__8b4751b6', 'check_video_audio_info__c6746b3e', 'check_mp3_audio_extraction__5e12822f6a68285e2467088dd7d25598', 'check_vlc_qt_autoload_extensions__8d3419de', 'check_vlc_time_advanced__2bf4fa278b8d46abde1b1e5c8949b96a', 'check_video_duration_saved__da7e9a1f', 'check_audio_conversion__70ff6e4ac087c63cd4c29ec50188d44d', 'check_vlc_video_on_top__81895b32', 'check_vlc_qt_privacy_ask__f900ec1c', 'check_qt_updates_notif__b84bb567', 'check_vlc_continue_playback_enabled', 'check_qt_continue_playback__90a64484', 'check_vlc_key_fullscreen__a57dec40610279f297f0f0fecc93f9ff', 'check_vlc_video_on_top__901155ec', 'check_is_video_file__4939fb90', 'check_snapshot_created__4e252238', 'check_audio_format__5d993657', 'check_snapshot_size__313dd5e1', 'check_audio_file_format__6475bf6e7aa0ef4599a6c11ec95f5406', 'check_vlc_start_paused__443cb77a', 'check_vlc_qt_notification__efa02ee6', 'check_vlc_video_transformed__93a4e125', 'check_vlc_playlist_count__fc73aaff', 'check_vlc_video_transformed__fc9f20bb', 'check_snapshot_saved__0706c584', 'check_audio_extracted__d0f84a157daee05fd0a57bb243a6ea5c', 'check_snapshot_created__1ce093b98b48496fecfb155d31ed7704', 'check_vlc_playback_rate__d5f08d8dc19ae133ba902b1660141393', 'check_vlc_random_mode__de28db3c', 'check_snapshot_directory__008263c0', 'check_qt_continue_playback__e7d03ebf', 'check_vlc_qt_fs_controller__dd2b0e06', 'check_video_rotation__ad793600129eb921cf29c68b343191af', 'check_vlc_bgcone__ed0d0c08edef23629254f099c29e5e89', 'check_vlc_screenshot_path__6c9e8f1d4a3b2e5f7890abcd12345678', 'check_video_bitrate__e0916c40', 'check_vlc_start_paused__a9847f6d', 'check_vlc_paused__9bec8597', 'check_vlc_global_hotkey_stop__dfe711eae38b1244e43a6b27c52f3be0', 'check_qt_system_tray__629fd64a', 'check_vlc_screenshot_format__bb7a94b9', 'check_snapshot_prefix__bdf11a45', 'check_vlc_loop_status__9e0433d6', 'check_snapshot_path__f6cd7aaf', 'check_video_duration__ff469bc6', 'check_vlc_always_on_top__b1199d7f2afb5cd10588e7e1e6cdc076', 'check_vlc_video_transformed__5e99ab87', 'check_qt_video_autoresize__6947249d', 'check_vlc_muted__95775285be787fb106c6c57a61517d2a', 'check_snapshot_format__a060a2ea', 'check_vlc_always_on_top__245c3f85', 'check_vlc_screenshot_saved__27395ab4949ce7b3fcdc03fed506b8b7', 'check_snapshot_count__8588412f', 'check_wav_conversion__580377a3f955e45d41d67d84cdd5fa88', 'check_vlc_minimal_view__99af05cec0829d789ac2d7bf7abe0481', 'check_vlc_loop_mode__6d45f22f', 'check_vlc_playback_rate__a884bc65', 'check_vlc_playing_file__24795e62', 'check_vlc_snapshot_directory__b9846f28', 'check_vlc_stopped__b038ba5d0dbf1c0974f6a77a9290fde5', 'check_vlc_video_transformed__bdd54294', 'check_vlc_repeat_mode__e611293b', 'check_video_metadata_extracted__bc3a25a7baab15992708f46f1d51e584', 'check_vlc_max_volume__dea0bbbbee03e4923af38de1331e256f', 'check_vlc_loop__869e58e4ccfbeafe3d06c94499c355f5', 'check_vlc_loop_mode__03941735', 'check_wallpaper_is_snapshot__b36f4827', 'check_srt_file_and_video__ed96ceb6', 'check_vlc_muted__eef15ca1487d76cdd0b6405615ef653d', 'check_video_framerate__8c156dd1', 'check_vlc_audio_muted__b5c6209b', 'check_vlc_qt_video_autoresize__06c10f73', 'check_vlc_start_paused__ae1f301e', 'check_vlc_loop_setting__0435b04a', 'check_audio_duration__ee9c4304f47d297d44dedaad1e2983d6', 'check_vlc_loop_mode__bd908fa8', 'check_audio_file_created__5d993657', 'check_vlc_aspect_ratio__8a47ea01', 'check_vlc_config_and_mp3_file__8f080098', 'check_vlc_metadata_network_access__e85db42df496eec95b9a79d8fc8ac5e4', 'check_vlc_subtitle_autoload__aec93e6e', 'check_qt_privacy_ask__f3480833', 'check_vlc_volume__c48d7656', 'check_video_codec__05e4689a', 'check_video_resolution__89b2d435', 'check_vlc_video_transformed__6e7fd09f', 'check_video_snapshot__1391f174']

def check_qt_fullscreen_control__f9958b0e(actual_config_path, expected, **options):
    """
    Checks if VLC's fullscreen controller setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_fullscreen_control" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_fullscreen_control = expected['expected_qt_fullscreen_control']
    if isinstance(expected_qt_fullscreen_control, int):
        expected_qt_fullscreen_control = str(expected_qt_fullscreen_control)
    try:
        qt_fullscreen_control = '1'
        for line in config_file.split('\n'):
            if 'qt-fs-controller=' in line:
                qt_fullscreen_control = line.split('=')[-1].strip()
        if qt_fullscreen_control == expected_qt_fullscreen_control:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_fullscreen_state__d7a7afc9(result, expected, **options):
    """
    Check if VLC is in fullscreen mode as expected.

    Args:
        result: Boolean indicating fullscreen state
        expected: Expected configuration
        **options: Additional options

    Returns:
        float: 1.0 if fullscreen state matches expected, 0.0 otherwise
    """
    should_be_fullscreen = expected.get('should_be_fullscreen', True)
    if should_be_fullscreen:
        return 1.0 if result else 0.0
    else:
        return 0.0 if result else 1.0

def check_vlc_paused__60670fd49b57e833635bf052410c24fc(actual_status_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC is currently paused on a specific file.
    """
    with open(actual_status_path, 'rb') as file:
        actual_status = file.read().decode('utf-8')
    tree = ElementTree.fromstring(actual_status)
    status = tree.find('state').text
    logger.info(f'VLC Status: {status}')
    if status == 'paused':
        if rule['type'] == 'file_name':
            file_paths = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="uri"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="name"]']
            file_info = None
            for path in file_paths:
                element = tree.find(path)
                if element is not None and element.text:
                    file_info = element.text
                    break
            if file_info:
                expected_filename = rule['file_name']
                import os
                actual_basename = os.path.basename(file_info)
                if actual_basename == expected_filename:
                    return 1
                if file_info.endswith(expected_filename):
                    return 1
                if expected_filename in file_info:
                    if file_info.endswith('/' + expected_filename) or file_info.endswith('\\' + expected_filename):
                        return 1
                logger.warning(f'File name mismatch - Expected: {expected_filename}, Found: {file_info}')
                return 0
            else:
                logger.warning(f'Could not find file information in VLC status XML for rule: {rule}')
                return 0
        else:
            logger.error(f"Unknown type: {rule['type']}")
            return 0
    else:
        logger.warning(f'VLC is not paused. Current state: {status}')
        return 0

def check_vlc_playback_state__fc6f4242(result, expected, **options):
    """
    Compare VLC playback state against expected value.

    Args:
        result: Current playback state from getter (str or None)
        expected: Expected playback state (str, e.g., 'paused', 'playing', 'stopped')
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC playback state result is None')
        return 0.0
    result_normalized = result.strip().lower()
    expected_normalized = expected.strip().lower()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        logger.info(f'Playback state mismatch - Expected: {expected}, Got: {result}')
        return 0.0

def check_vlc_loop__9eea6b3721652abbf42884414d2268c7(result, expected, **options):
    """
    Checks if VLC loop/repeat mode matches the expected state.

    Args:
        result: The current loop status from VLC (boolean)
        expected: The expected loop state from rules

    Returns:
        float: 1.0 if loop state matches, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC loop status is None')
        return 0.0
    expected_loop = expected.get('loop_enabled', True)
    if result == expected_loop:
        logger.info(f'VLC loop state matches expected: {expected_loop}')
        return 1.0
    else:
        logger.warning(f'VLC loop state mismatch - Expected: {expected_loop}, Got: {result}')
        return 0.0

def check_vlc_repeat_mode__9d232fb9(actual_config_path, expected):
    """
    Checks if VLC's repeat mode is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_repeat = expected.get('repeat_mode')
    if isinstance(expected_repeat, int):
        expected_repeat = str(expected_repeat)
    try:
        repeat_mode = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'repeat=' in line:
                repeat_mode = line.split('=')[-1].strip()
                break
        if repeat_mode == expected_repeat:
            return 1.0
        else:
            logger.warning(f'Repeat mode mismatch - Expected: {expected_repeat}, Found: {repeat_mode}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_qt_startminimized__5a4867ff(actual_config_path, expected):
    """Check VLC start minimized setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_startminimized = expected.get('expected_startminimized')
    if isinstance(expected_startminimized, int):
        expected_startminimized = str(expected_startminimized)
    try:
        startminimized = '0'
        for line in config_file.split('\n'):
            if 'qt-startminimized=' in line and (not line.strip().startswith('#')):
                startminimized = line.split('=')[-1].strip()
        if startminimized == expected_startminimized:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_subtitle_autodetect__d8ae66c5(actual_config_path, expected):
    """
    Checks if VLC's subtitle autodetect setting is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected configuration (from rules dict)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_subtitle_autodetect = expected['expected_subtitle_autodetect']
    if isinstance(expected_subtitle_autodetect, int):
        expected_subtitle_autodetect = str(expected_subtitle_autodetect)
    try:
        subtitle_autodetect = '1'
        for line in config_file.split('\n'):
            if 'sub-autodetect-file=' in line:
                subtitle_autodetect = line.split('=')[-1].strip()
        if subtitle_autodetect == expected_subtitle_autodetect:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_skip_frames__99a2f14e8a82b525ea717f160309e108(actual_config_path, expected):
    """
    Checks if VLC's skip frames setting is enabled/disabled as expected.
    """
    try:
        with open(actual_config_path, 'rb') as file:
            config_file = file.read().decode('utf-8')
    except UnicodeDecodeError:
        with open(actual_config_path, 'rb') as file:
            config_file = file.read().decode('latin-1')
    expected_skip_frames = expected['expected_skip_frames']
    if isinstance(expected_skip_frames, int):
        expected_skip_frames = str(expected_skip_frames)
    try:
        skip_frames = '0'
        for line in config_file.split('\n'):
            if 'skip-frames=' in line:
                skip_frames = line.split('=')[-1].strip()
        if skip_frames == expected_skip_frames:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_screenshot_format__5940c66b(actual_config_path, rule):
    """
    Checks if VLC's screenshot format setting is configured correctly.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_format = rule['expected_screenshot_format']
    try:
        screenshot_format = 'png'
        for line in config_file.split('\n'):
            if 'snapshot-format=' in line:
                screenshot_format = line.split('=')[-1].strip()
        if screenshot_format == expected_format:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_autosave_volume__9190a44d(actual_config_path, expected, **options):
    """
    Checks if VLC's autosave volume setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_autosave_volume" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_autosave_volume = expected['expected_qt_autosave_volume']
    if isinstance(expected_qt_autosave_volume, int):
        expected_qt_autosave_volume = str(expected_qt_autosave_volume)
    try:
        qt_autosave_volume = '1'
        for line in config_file.split('\n'):
            if 'qt-autosave-volume=' in line:
                qt_autosave_volume = line.split('=')[-1].strip()
        if qt_autosave_volume == expected_qt_autosave_volume:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_loop_mode__23c7a390(actual_config_path, expected):
    """
    Checks if VLC's loop/repeat mode is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected configuration (from rules dict)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_loop_mode = expected['expected_loop_mode']
    if isinstance(expected_loop_mode, int):
        expected_loop_mode = str(expected_loop_mode)
    try:
        loop_mode = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'loop=' in line and 'loop-filter=' not in line:
                loop_mode = line.split('=')[-1].strip()
        if loop_mode == expected_loop_mode:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_video_effects__84f74170(result, expected, **options):
    """
    Check if mirror/horizontal flip video effect is enabled in VLC configuration.

    Args:
        result: Dictionary of VLC video effect settings with 'mirror_enabled' key
        expected: Expected configuration with 'mirror_enabled' to check
        **options: Additional options

    Returns:
        float: 1.0 if mirror effect is enabled as expected, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    expected_mirror = expected.get('mirror_enabled', True)
    actual_mirror = result.get('mirror_enabled', False)
    if expected_mirror == actual_mirror:
        return 1.0
    return 0.0

def check_vlc_playlist_tree__2b678efeda2f5b06f01e1ca5035b98c7(actual_config_path, rule):
    """
    Checks if VLC's playlist tree view is enabled/disabled as expected.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_playlist_tree = rule['expected_playlist_tree']
    if isinstance(expected_playlist_tree, int):
        expected_playlist_tree = str(expected_playlist_tree)
    try:
        qt_pl_showflags = '1'
        for line in config_file.split('\n'):
            if 'qt-pl-showflags=' in line:
                qt_pl_showflags = line.split('=')[-1].strip()
        if qt_pl_showflags == expected_playlist_tree:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_always_on_top__850d1898(actual_config_path, expected):
    """
    Checks if VLC's always-on-top setting is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected configuration (from rules dict)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_always_on_top = expected['expected_always_on_top']
    if isinstance(expected_always_on_top, int):
        expected_always_on_top = str(expected_always_on_top)
    try:
        always_on_top = '0'
        for line in config_file.split('\n'):
            if 'video-on-top=' in line:
                always_on_top = line.split('=')[-1].strip()
        if always_on_top == expected_always_on_top:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_snapshot_format__772d6c83(actual_config_path, expected):
    """Check VLC snapshot format setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_format = expected.get('expected_snapshot_format')
    try:
        snapshot_format = 'png'
        for line in config_file.split('\n'):
            if 'snapshot-format=' in line and (not line.strip().startswith('#')):
                snapshot_format = line.split('=')[-1].strip()
        if snapshot_format == expected_format:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_bookmark__1c6b323b(result_state, expected, **options):
    """
    Checks if a VLC bookmark exists at the expected position by verifying the accessibility tree.

    Args:
        result_state: Dict containing bookmark dialog information from getter
        expected: Dict with 'has_bookmark' (bool) and 'expected_position' (float, in seconds)
        **options: Additional options

    Returns:
        float: 1.0 if bookmark found at correct position, 0.0 otherwise
    """
    if not isinstance(result_state, dict):
        logger.error(f'Invalid result_state type: {type(result_state)}, expected dict')
        return 0.0
    should_have_bookmark = expected.get('has_bookmark', True)
    expected_position = expected.get('expected_position', None)
    has_dialog = result_state.get('has_dialog', False)
    if not has_dialog:
        logger.warning('Bookmark dialog not found in accessibility tree')
        return 0.0 if should_have_bookmark else 1.0
    bookmark_entries = result_state.get('bookmark_entries', [])
    if not should_have_bookmark:
        return 0.0 if len(bookmark_entries) > 0 else 1.0
    if len(bookmark_entries) == 0:
        logger.warning('No bookmark entries found in dialog')
        return 0.0
    if expected_position is None:
        logger.info(f'Found {len(bookmark_entries)} bookmark(s): {bookmark_entries}')
        return 1.0

    def time_matches_position(time_str, position_seconds, tolerance=2.0):
        """Check if a time string matches the expected position within tolerance."""
        import re
        mm_ss_match = re.match('(\\d{1,2}):(\\d{2})(?:\\.(\\d+))?', time_str)
        if mm_ss_match:
            minutes = int(mm_ss_match.group(1))
            seconds = int(mm_ss_match.group(2))
            milliseconds = float('0.' + mm_ss_match.group(3)) if mm_ss_match.group(3) else 0.0
            total_seconds = minutes * 60 + seconds + milliseconds
            return abs(total_seconds - position_seconds) <= tolerance
        sec_match = re.match('(\\d+(?:\\.\\d+)?)\\s*s?', time_str)
        if sec_match:
            total_seconds = float(sec_match.group(1))
            return abs(total_seconds - position_seconds) <= tolerance
        return False
    for entry in bookmark_entries:
        if time_matches_position(entry, expected_position):
            logger.info(f'Found matching bookmark at position {entry} (expected ~{expected_position}s)')
            return 1.0
    logger.warning(f'No bookmark found at expected position {expected_position}s. Found: {bookmark_entries}')
    return 0.0

def check_vlc_stopped__c96690ce(actual_status_path: str, expected) -> float:
    """
    Checks if VLC video is stopped.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Expected state ('stopped')

    Returns:
        1.0 if VLC is stopped, 0.0 otherwise
    """
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        status = tree.find('state').text
        logger.info(f'VLC Status: {status}')
        expected_state = expected if isinstance(expected, str) else expected.get('state', 'stopped')
        if status == expected_state:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC stopped status: {e}')
        return 0.0

def check_vlc_recent_media__aa13955b1d571b5f9f91c4b58d07755a(result, expected, **options):
    """Check if VLC recent media contains expected video.

    Args:
        result: Dict from getter with 'media_items' list
        expected: Dict with 'contains_video' string from rules
        **options: Additional options

    Returns:
        float: 1.0 if video found in recent media, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Invalid result format: {result}')
        return 0.0
    media_items = result.get('media_items', [])
    raw_content = result.get('raw_content', '')
    expected_video = expected.get('contains_video', '')
    if not expected_video:
        logger.error('No expected_video specified in rules')
        return 0.0
    for item in media_items:
        if expected_video in item:
            logger.info(f'Found {expected_video} in media items')
            return 1.0
    if expected_video in raw_content:
        logger.info(f'Found {expected_video} in VLC config raw content')
        return 1.0
    logger.warning(f'Video {expected_video} not found in VLC recent media')
    return 0.0

def check_vlc_video_transformed__78c019ea(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly rotated 180 degrees.

    For 180° rotation verification:
    - Dimensions remain the same (not swapped like 90°/270° rotations)
    - Extracts frames from both source and output videos
    - Rotates source frame 180° and compares with output frame using pixel comparison
    - This ensures the specific 180° rotation was applied, not just any transformation

    Args:
        result: Combined output from ls, stat, ffprobe, and frame comparison commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was rotated 180°, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    remaining_frame_data = probe_parts[1] if len(probe_parts) > 1 else ''
    frame_parts = remaining_frame_data.split('---FRAME_EXTRACT---')
    source_video_info = frame_parts[0].strip() if len(frame_parts) > 0 else ''
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    if len(frame_parts) > 1:
        frame_data = frame_parts[1]
        pixel_compare_parts = frame_data.split('---PIXEL_COMPARE---')
        if len(pixel_compare_parts) > 1:
            pixel_compare_output = pixel_compare_parts[1].strip()
            if 'ERROR' in pixel_compare_output:
                logger.warning('Frame comparison failed - falling back to basic checks')
            elif 'COMPARE_DONE' in pixel_compare_output or pixel_compare_output:
                rmse_match = re.search('([\\d.]+)\\s+\\(([\\d.]+)\\)', pixel_compare_output)
                if rmse_match:
                    rmse_normalized = float(rmse_match.group(2))
                    RMSE_THRESHOLD = 0.05
                    if rmse_normalized <= RMSE_THRESHOLD:
                        logger.info(f'180° rotation verified via frame comparison (RMSE: {rmse_normalized:.6f})')
                        return 1.0
                    else:
                        logger.warning(f'Frame comparison shows rotation mismatch (RMSE: {rmse_normalized:.6f} > {RMSE_THRESHOLD})')
                        return 0.0
                else:
                    logger.info('Could not parse RMSE from pixel comparison')
            hash_parts = pixel_compare_parts[0].split('---ROTATED_HASH---')
            if len(hash_parts) > 1:
                hash_data = hash_parts[1]
                output_hash_parts = hash_data.split('---OUTPUT_HASH---')
                if len(output_hash_parts) > 1:
                    rotated_hash_str = output_hash_parts[0].strip()
                    output_hash_str = output_hash_parts[1].strip()
                    try:
                        rotated_hash = float(rotated_hash_str)
                        output_hash = float(output_hash_str)
                        hash_diff = abs(rotated_hash - output_hash)
                        HASH_THRESHOLD = 0.02
                        if hash_diff <= HASH_THRESHOLD:
                            logger.info(f'180° rotation verified via perceptual hash (diff: {hash_diff:.6f})')
                            return 1.0
                        else:
                            logger.warning(f'Perceptual hash mismatch (diff: {hash_diff:.6f} > {HASH_THRESHOLD})')
                    except (ValueError, TypeError) as e:
                        logger.info(f'Could not parse hash values: {e}')
    logger.warning('Frame-level verification unavailable - checking for basic transformation evidence')
    if source_info:
        source_width = source_info.get('width', '')
        source_height = source_info.get('height', '')
        output_width = output_info.get('width', '')
        output_height = output_info.get('height', '')
        if source_width and source_height and output_width and output_height:
            if source_width != output_width or source_height != output_height:
                logger.warning('Dimensions changed - this suggests 90°/270° rotation or cropping, not 180°')
                return 0.0
    output_filesize_str = output_filesize_str.strip()
    source_filesize_str = source_filesize_str.strip()
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        output_filesize = None
        source_filesize = None
    if output_filesize is not None and source_filesize is not None:
        if output_filesize == source_filesize:
            logger.warning('File sizes identical - likely just copied/renamed without transformation')
            return 0.0
        else:
            size_diff_percent = abs(output_filesize - source_filesize) / max(source_filesize, 1) * 100
            logger.info(f'File size changed by {size_diff_percent:.2f}% - some transformation occurred')
            return 0.3
    logger.warning('Could not verify 180° rotation - insufficient evidence')
    return 0.0

def check_qt_playlist_tree__bf5eef29(actual_config_path, expected):
    """
    Check if VLC playlist is configured to show in tree view.
    This checks the qt-playlist-tree setting in VLC config.

    Args:
        actual_config_path: Path to vlcrc config file
        expected: Expected value (0=list view, 1=tree view)

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_playlist_tree = expected.get('expected_qt_playlist_tree')
    if isinstance(expected_qt_playlist_tree, int):
        expected_qt_playlist_tree = str(expected_qt_playlist_tree)
    try:
        qt_playlist_tree = '1'
        for line in config_file.split('\n'):
            if 'qt-playlist-tree=' in line:
                qt_playlist_tree = line.split('=')[-1].strip()
        if qt_playlist_tree == expected_qt_playlist_tree:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_snapshot_exists__57da68f2(result, expected, **options):
    """
    Check if a VLC snapshot file exists.

    This metric verifies that a snapshot was taken during task execution.
    The task config includes a preconfig step that removes any existing
    vlcsnap-*.png files before task execution, ensuring that only snapshots
    created during the task will be detected by the getter.

    Args:
        result: Filename returned from getter (str if snapshot found, None otherwise)
        expected: Expected state (dict with 'snapshot_should_exist' boolean)
        **options: Additional options

    Returns:
        float: 1.0 if snapshot exists as expected, 0.0 otherwise
    """
    snapshot_should_exist = expected.get('snapshot_should_exist', True)
    if snapshot_should_exist:
        if result and isinstance(result, str) and result.startswith('vlcsnap-') and result.endswith('.png'):
            return 1.0
        return 0.0
    else:
        if result is None:
            return 1.0
        return 0.0

def check_mp3_file_valid__ce2e74ee9437f284464037da6df4e453(result, expected, **options):
    """
    Check if MP3 file exists, is valid audio, and meets size requirements.

    Args:
        result: Dict from getter with file info
        expected: Dict with rules for validation
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        logger.warning('MP3 file does not exist')
        return 0.0
    if result.get('is_audio', False):
        score += 0.4
    else:
        logger.warning('File exists but is not a valid audio file')
    min_size = expected.get('min_size', 100000)
    file_size = result.get('size', 0)
    if file_size >= min_size:
        score += 0.2
    else:
        logger.warning(f'File size {file_size} bytes is below minimum {min_size} bytes')
    logger.info(f"MP3 validation score: {score} (exists: {result.get('exists')}, is_audio: {result.get('is_audio')}, size: {file_size})")
    return score

def check_vlc_subtitle_folder__73747be5(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's subtitle autoload folder is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_subtitle_path = rule['subtitle_folder_path']
    try:
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'sub-autodetect-path' in line:
                current_path = line.split('=')[-1].strip()
                if current_path == expected_subtitle_path:
                    return 1.0
                else:
                    return 0.0
        return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_screenshot_path__fe752f70(actual_config_path, expected):
    """
    Checks if VLC's screenshot path is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_path = expected.get('screenshot_path')
    try:
        screenshot_path = None
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-path=' in line:
                screenshot_path = line.split('=')[-1].strip()
                break
        if screenshot_path == expected_path:
            return 1.0
        else:
            logger.warning(f'Screenshot path mismatch - Expected: {expected_path}, Found: {screenshot_path}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_paused__8f5f9872b32c7843d08f183a299b915a(actual_status_path: str, expected: Dict[str, str]) -> float:
    """
    Checks if VLC is in paused state.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Dict with expected state (should contain expected_state='paused')

    Returns:
        1.0 if VLC is paused, 0.0 otherwise
    """
    if not actual_status_path:
        logger.error('No VLC status path provided')
        return 0.0
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        status = tree.find('state').text
        expected_state = expected.get('expected_state', 'paused')
        logger.info(f'VLC Status: {status}, Expected: {expected_state}')
        if status == expected_state:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC paused state: {e}')
        return 0.0

def check_vlc_screenshot_exists__167faa79(result, expected, **options):
    """
    Check if a VLC screenshot filename was found.

    Args:
        result: Filename string returned by getter (empty if not found)
        expected: Dict with 'rules' containing 'filename_pattern' to match
        **options: Additional options

    Returns:
        float: 1.0 if filename matches pattern, 0.0 otherwise
    """
    filename_pattern = expected.get('filename_pattern', 'vlc-snap')
    if not result:
        logger.warning('No screenshot filename found')
        return 0.0
    if filename_pattern in result:
        logger.info(f'Screenshot verified: {result}')
        return 1.0
    else:
        logger.warning(f"Filename pattern '{filename_pattern}' not in '{result}'")
        return 0.0

def check_ogg_audio_extraction__c03bf6861f4e39f709595e582520ac0f(result, expected, **options):
    """
    Verify OGG Vorbis audio file was extracted from video.
    Progressive scoring for file creation, naming, format, and size.
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.35
        logger.info(f"OGG file exists: {result.get('filename', 'unknown')}")
    else:
        logger.warning('OGG file not created')
        return score
    expected_filename = expected.get('filename', '')
    if result.get('filename', '') == expected_filename:
        score += 0.3
        logger.info(f'Filename matches: {expected_filename}')
    else:
        logger.warning(f"Filename mismatch - Expected: {expected_filename}, Got: {result.get('filename', '')}")
    if result.get('is_ogg', False):
        score += 0.2
        logger.info(f"OGG format verified: {result.get('mime_type', 'unknown')}")
    else:
        logger.warning(f"Not OGG format: {result.get('mime_type', 'unknown')}")
    min_size = expected.get('min_size', 80000)
    if result.get('size', 0) >= min_size:
        score += 0.15
        logger.info(f"File size OK: {result.get('size', 0)} bytes (min: {min_size})")
    else:
        logger.warning(f"File size insufficient: {result.get('size', 0)} bytes (min: {min_size})")
    return score

def check_qt_autosave_volume__f8ad1f01(actual_config_path, rule):
    """Check if VLC autosave volume setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_autosave_volume = rule['expected_qt_autosave_volume']
    if isinstance(expected_qt_autosave_volume, int):
        expected_qt_autosave_volume = str(expected_qt_autosave_volume)
    try:
        qt_autosave_volume = '0'
        for line in config_file.split('\n'):
            if 'qt-autosave-volume=' in line:
                qt_autosave_volume = line.split('=')[-1].strip()
        if qt_autosave_volume == expected_qt_autosave_volume:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_qt_video_autoresize__a74697ad(actual_config_path, expected, **options):
    """
    Checks if VLC's video auto-resize setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_video_autoresize" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_video_autoresize = expected['expected_qt_video_autoresize']
    if isinstance(expected_qt_video_autoresize, int):
        expected_qt_video_autoresize = str(expected_qt_video_autoresize)
    try:
        qt_video_autoresize = '1'
        for line in config_file.split('\n'):
            if 'qt-video-autoresize=' in line:
                qt_video_autoresize = line.split('=')[-1].strip()
        if qt_video_autoresize == expected_qt_video_autoresize:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_on_top__e48d3a20(actual_config_path, rule):
    """
    Checks if VLC's video-on-top setting is set to the expected value.
    video-on-top=1 means always on top is enabled.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_video_on_top = rule['expected_video_on_top']
    if isinstance(expected_video_on_top, int):
        expected_video_on_top = str(expected_video_on_top)
    try:
        video_on_top = '0'
        for line in config_file.split('\n'):
            if 'video-on-top=' in line:
                video_on_top = line.split('=')[-1].strip()
        if video_on_top == expected_video_on_top:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_mp3_count__be089f93(result, expected, **options):
    """Check if MP3 file count matches expected value.

    Args:
        result: MP3 file count from getter
        expected: Expected MP3 file count
        **options: Additional options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    if result is None:
        return 0.0
    if result == expected:
        return 1.0
    else:
        return 0.0

def check_vlc_repeat_loop_enabled(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's repeat or loop mode is enabled in the configuration file.

    VLC uses two settings in vlcrc:
    - loop=1: Repeat all (loop through playlist)
    - repeat=1: Repeat current item

    Args:
        actual_config_path: Path to the VLC config file
        rule: Dictionary with 'expected_repeat_or_loop' key (0 or 1)

    Returns:
        1.0 if repeat or loop is enabled as expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_enabled = rule.get('expected_repeat_or_loop', '1')
    if isinstance(expected_enabled, int):
        expected_enabled = str(expected_enabled)
    try:
        loop_enabled = '0'
        repeat_enabled = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if line.startswith('loop='):
                loop_enabled = line.split('=')[-1].strip()
            elif line.startswith('repeat='):
                repeat_enabled = line.split('=')[-1].strip()
        actual_enabled = '1' if loop_enabled == '1' or repeat_enabled == '1' else '0'
        logger.info(f'VLC repeat/loop status - loop={loop_enabled}, repeat={repeat_enabled}, combined={actual_enabled}')
        if actual_enabled == expected_enabled:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_video_transformed__a899663b(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly rotated 90 degrees clockwise.

    For 90-degree rotation, we verify:
    1. Output file exists with correct name
    2. Output is a valid video file
    3. Dimensions are swapped (width becomes height, height becomes width)
    4. File was re-encoded (different file size)

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was rotated 90 degrees, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    if not source_info:
        logger.warning('Cannot verify rotation without source video info')
        return 0.5
    output_width = output_info.get('width', '')
    output_height = output_info.get('height', '')
    source_width = source_info.get('width', '')
    source_height = source_info.get('height', '')
    if not (output_width and output_height and source_width and source_height):
        logger.warning('Missing dimension information')
        return 0.0
    dimensions_swapped = output_width == source_height and output_height == source_width
    if not dimensions_swapped:
        logger.warning(f'Dimensions NOT swapped for 90-degree rotation. Source: {source_width}x{source_height}, Output: {output_width}x{output_height}. Expected output to be {source_height}x{source_width}.')
        return 0.0
    logger.info(f'Dimensions correctly swapped: {source_width}x{source_height} -> {output_width}x{output_height}')
    if output_filesize is not None and source_filesize is not None:
        if output_filesize == source_filesize:
            logger.warning('File sizes identical - likely renamed/copied, not rotated')
            return 0.5
        else:
            logger.info(f'File was re-encoded: {source_filesize} -> {output_filesize} bytes')
    logger.info('Video rotation (90 degrees) verified successfully')
    return 1.0

def check_qt_video_autoresize__b7688244(actual_config_path, rule):
    """Check if VLC video auto-resize setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_video_autoresize = rule['expected_qt_video_autoresize']
    if isinstance(expected_qt_video_autoresize, int):
        expected_qt_video_autoresize = str(expected_qt_video_autoresize)
    try:
        qt_video_autoresize = '1'
        for line in config_file.split('\n'):
            if 'qt-video-autoresize=' in line:
                qt_video_autoresize = line.split('=')[-1].strip()
        if qt_video_autoresize == expected_qt_video_autoresize:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_audio_output__c4a2eaec(actual_config_path, expected):
    """
    Checks if VLC's audio output module is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_output = expected.get('audio_output')
    try:
        audio_output = None
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'aout=' in line:
                audio_output = line.split('=')[-1].strip()
                break
        if audio_output == expected_output:
            return 1.0
        else:
            logger.warning(f'Audio output mismatch - Expected: {expected_output}, Found: {audio_output}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_recordings_folder__055ce0bbe52cb194135b41b067eca986(actual_config_path: str, expected: Dict) -> float:
    """
    Checks if VLC's recording folder is set to the expected value.

    Args:
        actual_config_path: Path to the VLC config file
        expected: Dict containing 'recording_folder_path' key with expected path

    Returns:
        1.0 if recording folder matches expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_recording_path = expected['recording_folder_path']
    try:
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'input-record-path' in line:
                current_path = line.split('=')[-1].strip()
                if current_path == expected_recording_path:
                    return 1.0
                else:
                    logger.warning(f'Recording path mismatch - Expected: {expected_recording_path}, Found: {current_path}')
                    return 0.0
        logger.warning('input-record-path not found in VLC config')
        return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_updates_notif__f6b51378(actual_config_path, expected, **options):
    """
    Checks if VLC's update notification setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_updates_notif" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_updates_notif = expected['expected_qt_updates_notif']
    if isinstance(expected_qt_updates_notif, int):
        expected_qt_updates_notif = str(expected_qt_updates_notif)
    try:
        qt_updates_notif = '1'
        for line in config_file.split('\n'):
            if 'qt-updates-notif=' in line:
                qt_updates_notif = line.split('=')[-1].strip()
        if qt_updates_notif == expected_qt_updates_notif:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_volume__9cc567291e2f7ec7556a7b0f92f9b42b(vlc_status: Optional[Dict], rule: Dict) -> float:
    """
    Checks if VLC is playing media and the volume is set to the expected level.
    Verifies both playback state and volume with tolerance.
    Returns partial credit: 0.5 for playing state, 0.5 for correct volume.
    """
    if vlc_status is None:
        logger.warning('Could not retrieve VLC status')
        return 0.0
    score = 0.0
    state = vlc_status.get('state')
    expected_filename = rule.get('expected_filename', 'Rick Astley - Never Gonna Give You Up')
    if state == 'playing':
        score += 0.5
        logger.info(f'VLC is playing (state: {state})')
        filename = vlc_status.get('filename', '')
        if filename and expected_filename.lower() in filename.lower():
            logger.info(f'Correct media file is playing: {filename}')
        elif filename:
            logger.info(f'Media is playing but may not match expected: {filename}')
    else:
        logger.warning(f'VLC is not playing (state: {state})')
    actual_volume = vlc_status.get('volume')
    if actual_volume is not None:
        expected_volume = rule.get('expected_volume', 128)
        tolerance = rule.get('tolerance', 15)
        logger.info(f'VLC Volume - Expected: {expected_volume}, Actual: {actual_volume}, Tolerance: {tolerance}')
        if abs(actual_volume - expected_volume) <= tolerance:
            score += 0.5
            logger.info(f'Volume is correct: {actual_volume}')
        else:
            logger.warning(f'Volume mismatch - Expected: {expected_volume} (±{tolerance}), Actual: {actual_volume}')
    else:
        logger.warning('Could not retrieve VLC volume')
    logger.info(f'Final score: {score}')
    return score

def check_snapshot_format__ccaf0dcd(actual_config_path, expected):
    """
    Check if VLC video snapshot format is set to the expected format.
    This checks the snapshot-format setting in VLC config.

    Args:
        actual_config_path: Path to vlcrc config file
        expected: Expected format (e.g., 'png', 'jpg', 'tiff')

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_snapshot_format = expected.get('expected_snapshot_format')
    if expected_snapshot_format:
        expected_snapshot_format = expected_snapshot_format.lower()
    try:
        snapshot_format = 'png'
        for line in config_file.split('\n'):
            if 'snapshot-format=' in line:
                snapshot_format = line.split('=')[-1].strip().lower()
        if snapshot_format == expected_snapshot_format:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_bookmarks_created__18214e12ab9f5ae2ef34a7b82632e6f9(result, expected, **options):
    """
    Check if VLC bookmarks file was properly created.

    Args:
        result: Dict with bookmarks status from getter
        expected: Dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    if result.get('has_video_entry', False):
        score += 0.5
    return score

def check_vlc_playback_rate__c979999ccfa7b4560c7eb8656e505ff9(result, expected, **options):
    """
    Checks if VLC playback rate matches the expected rate.

    Args:
        result: The current playback rate from VLC (float)
        expected: The expected rate from rules

    Returns:
        float: 1.0 if rate matches (within tolerance), 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC playback rate is None')
        return 0.0
    expected_rate = expected.get('rate', 1.0)
    tolerance = expected.get('tolerance', 0.05)
    if abs(result - expected_rate) <= tolerance:
        logger.info(f'VLC playback rate matches expected: {expected_rate} (actual: {result})')
        return 1.0
    else:
        logger.warning(f'VLC playback rate mismatch - Expected: {expected_rate}, Got: {result}')
        return 0.0

def check_qt_privacy_ask__8eb13985(actual_config_path, rule):
    """Check if VLC privacy/network interaction setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_privacy_ask = rule['expected_qt_privacy_ask']
    if isinstance(expected_qt_privacy_ask, int):
        expected_qt_privacy_ask = str(expected_qt_privacy_ask)
    try:
        qt_privacy_ask = '1'
        for line in config_file.split('\n'):
            if 'qt-privacy-ask=' in line:
                qt_privacy_ask = line.split('=')[-1].strip()
        if qt_privacy_ask == expected_qt_privacy_ask:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_audio_extracted__dde7bdf2(result, expected, **options):
    """
    Check if audio was properly extracted from video.

    Args:
        result: Path to the extracted audio.mp3 file
        expected: Dict with "min_size_bytes" and "check_format" keys
        **options: Additional options

    Returns:
        float: 1.0 if audio extracted correctly, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        if not os.path.exists(result):
            logger.warning(f'File does not exist: {result}')
            return 0.0
        file_size = os.path.getsize(result)
        min_size = expected.get('min_size_bytes', 0)
        if file_size < min_size:
            logger.info(f'File size {file_size} is below minimum {min_size}')
            return 0.0
        if expected.get('check_format', False):
            with open(result, 'rb') as f:
                header = f.read(3)
                if not (header.startswith(b'ID3') or (header[0] == 255 and header[1] & 224 == 224)):
                    logger.warning(f'File does not appear to be valid MP3 format')
                    return 0.0
        logger.info(f'Audio file validated: {file_size} bytes')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking audio file: {e}')
        return 0.0

def check_qt_privacy_ask__132d5c78(actual_config_path, rule):
    """
    Checks if VLC's qt-privacy-ask setting is set to the expected value.
    qt-privacy-ask=0 means don't ask for privacy/network policy (already answered).
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_privacy_ask = rule['expected_qt_privacy_ask']
    if isinstance(expected_qt_privacy_ask, int):
        expected_qt_privacy_ask = str(expected_qt_privacy_ask)
    try:
        qt_privacy_ask = '1'
        for line in config_file.split('\n'):
            if 'qt-privacy-ask=' in line:
                qt_privacy_ask = line.split('=')[-1].strip()
        if qt_privacy_ask == expected_qt_privacy_ask:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def is_vlc_playing_youtube__bba3381f(actual_status_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC is currently playing a YouTube video.

    This metric handles YouTube URL transformation: VLC transforms YouTube URLs
    (youtube.com/watch?v=...) into actual streaming URLs (googlevideo.com domains)
    when playing. The metric verifies by:
    1. Checking VLC is in playing state
    2. For YouTube URLs: extracting and comparing video IDs from both expected URL
       and VLC status, OR checking for googlevideo.com domain
    3. Logging actual URL found for debugging

    Args:
        actual_status_path: Path to VLC status XML file
        rule: Dict containing either:
            - 'url': YouTube URL to verify (e.g., 'https://www.youtube.com/watch?v=9bZkp7q19f0')
            OR
            - 'domain': Expected domain (e.g., 'googlevideo.com')
              AND 'video_id': YouTube video ID (e.g., '9bZkp7q19f0')

    Returns:
        1.0 if VLC is playing the expected YouTube video, 0.0 otherwise
    """
    with open(actual_status_path, 'rb') as file:
        actual_status = file.read().decode('utf-8')
    tree = ElementTree.fromstring(actual_status)
    status = tree.find('state').text
    logger.info(f'VLC Status: {status}')
    if status != 'playing':
        logger.warning('VLC is not in playing state')
        return 0.0
    url_paths = ['information/category[@name="meta"]/info[@name="url"]', 'information/category[@name="meta"]/info[@name="URI"]', 'information/category[@name="meta"]/info[@name="location"]', 'information/category[@name="meta"]/info[@name="title"]', 'information/category[@name="meta"]/info[@name="filename"]']
    file_info = None
    for path in url_paths:
        element = tree.find(path)
        if element is not None and element.text:
            file_info = element.text
            logger.info(f"Found URL info at '{path}': {file_info}")
            break
    if not file_info:
        logger.warning('Could not find URL information in VLC status XML')
        return 0.0
    expected_url = rule.get('url', '')
    expected_domain = rule.get('domain', '')
    expected_video_id = rule.get('video_id', None)
    if expected_url and (not expected_video_id):
        expected_video_id = extract_youtube_video_id(expected_url)
        if expected_video_id:
            logger.info(f'Extracted video ID from expected URL: {expected_video_id}')
        else:
            logger.warning(f'Could not extract video ID from expected URL: {expected_url}')
    actual_video_id = extract_youtube_video_id(file_info)
    if actual_video_id:
        logger.info(f'Extracted video ID from VLC URL: {actual_video_id}')
    if expected_video_id and actual_video_id:
        if expected_video_id == actual_video_id:
            logger.info(f'Video ID match: {expected_video_id}')
            return 1.0
        else:
            logger.warning(f'Video ID mismatch - Expected: {expected_video_id}, Found: {actual_video_id}')
            return 0.0
    if expected_domain or expected_url:
        if 'youtube.com' in expected_url or 'youtu.be' in expected_url:
            expected_domain = 'googlevideo.com'
            logger.info(f'YouTube URL detected, accepting googlevideo.com domain')
        if expected_domain:
            if '://' in file_info:
                try:
                    parsed = urlparse(file_info)
                    if expected_domain in parsed.netloc or parsed.netloc.endswith(expected_domain):
                        logger.info(f'Domain match: {parsed.netloc} contains {expected_domain}')
                        return 1.0
                    else:
                        logger.warning(f'Domain mismatch - Expected domain: {expected_domain}, Found: {parsed.netloc}')
                        return 0.0
                except Exception as e:
                    logger.error(f'URL parsing error: {e}')
                    if expected_domain in file_info:
                        logger.info(f'Domain found in URL via substring match: {expected_domain}')
                        return 1.0
                    return 0.0
            else:
                if expected_domain in file_info:
                    logger.info(f'Domain found via substring match: {expected_domain}')
                    return 1.0
                logger.warning(f'Domain not found in file info: {file_info}')
                return 0.0
    logger.error('Could not verify YouTube video - no video ID match and no domain verification possible')
    return 0.0

def check_vlc_video_transformed__cb59ca2e(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly rotated to the correct orientation.

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename', 'source_filename', 'expected_rotation', and 'source_rotation' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and is rotated to correct orientation, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    expected_rotation = expected.get('expected_rotation', '0')
    source_rotation = expected.get('source_rotation', '180')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    output_rotation = output_info.get('rotate', '0')
    output_rotation = output_rotation if output_rotation else '0'
    expected_rotation = expected_rotation if expected_rotation else '0'
    if output_rotation != expected_rotation:
        logger.warning(f'Video rotation is {output_rotation}, expected {expected_rotation}')
        return 0.0
    logger.info(f'Video correctly rotated to {output_rotation} degrees')
    if source_info and output_filesize is not None and (source_filesize is not None):
        source_rotation_actual = source_info.get('rotate', '0')
        source_rotation_actual = source_rotation_actual if source_rotation_actual else '0'
        if source_rotation_actual != expected_rotation:
            if output_filesize == source_filesize:
                logger.warning('File sizes identical - likely renamed/copied, not transformed')
                return 0.0
            else:
                logger.info(f'Video re-encoded: size changed from {source_filesize} to {output_filesize}')
        else:
            logger.warning(f'Source video already had correct rotation {source_rotation_actual}')
    logger.info('Video rotation verified successfully')
    return 1.0

def check_vlc_paused__df488c66748635e65016c0e9e7cc92c2(result, expected, **options):
    """
    Checks if VLC is in the paused state.

    Args:
        result: The current playback state from VLC (string)
        expected: The expected state from rules (should be 'paused')

    Returns:
        float: 1.0 if paused, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC playback state is None')
        return 0.0
    expected_state = expected.get('state', 'paused')
    if result.lower() == expected_state.lower():
        logger.info(f'VLC state matches expected: {expected_state}')
        return 1.0
    else:
        logger.warning(f'VLC state mismatch - Expected: {expected_state}, Got: {result}')
        return 0.0

def check_vlc_screenshot_folder__3d380bd3(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's screenshot folder is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_screenshot_path = rule['screenshot_folder_path']
    try:
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-path' in line:
                current_path = line.split('=')[-1].strip()
                if current_path == expected_screenshot_path:
                    return 1.0
                else:
                    return 0.0
        return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_qt_pause_minimized__7c31b71a(actual_config_path, expected):
    """Check VLC pause when minimized setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_pause_minimized = expected.get('expected_pause_minimized')
    if isinstance(expected_pause_minimized, int):
        expected_pause_minimized = str(expected_pause_minimized)
    try:
        pause_minimized = '0'
        for line in config_file.split('\n'):
            if 'qt-pause-minimized=' in line and (not line.strip().startswith('#')):
                pause_minimized = line.split('=')[-1].strip()
        if pause_minimized == expected_pause_minimized:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_minimized__667c58f7(actual_window_info, expected):
    """
    Checks if VLC window is minimized and matches expected state.

    Args:
        actual_window_info: Window information dictionary
        expected: Expected state dictionary containing 'minimized' key

    Returns:
        1.0 if window state matches expected, 0.0 otherwise
    """
    try:
        expected_minimized = expected.get('minimized', False)
        is_minimized = actual_window_info is None or not actual_window_info or actual_window_info.get('width', 0) == 0 or (actual_window_info.get('height', 0) == 0)
        if is_minimized == expected_minimized:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC minimized status: {e}')
        return 0.0

def check_vlc_muted__cf575cad41cf5daac3aae885ce8cfe2a(actual_status: Optional[Dict], rule: Dict) -> float:
    """
    Checks if VLC's mute status and playback state match the expected state.
    Verifies both that audio is muted AND that a video file is playing.
    """
    if actual_status is None:
        logger.warning('Could not retrieve VLC status')
        return 0.0
    expected_muted = rule.get('expected_muted', True)
    expected_state = rule.get('expected_state', 'playing')
    actual_muted = actual_status.get('is_muted')
    actual_playback_state = actual_status.get('state')
    has_media = actual_status.get('has_media', False)
    logger.info(f'VLC Status Check - Expected Muted: {expected_muted}, Actual Muted: {actual_muted}')
    logger.info(f'VLC Status Check - Expected State: {expected_state}, Actual State: {actual_playback_state}')
    logger.info(f'VLC Status Check - Has Media: {has_media}')
    if actual_muted != expected_muted:
        logger.warning(f'Mute status mismatch - Expected: {expected_muted}, Actual: {actual_muted}')
        return 0.0
    if actual_playback_state != expected_state:
        logger.warning(f'Playback state mismatch - Expected: {expected_state}, Actual: {actual_playback_state}')
        return 0.0
    if not has_media:
        logger.warning('No media file loaded in VLC')
        return 0.0
    logger.info('VLC status matches all expected conditions')
    return 1.0

def check_vlc_snapshot_format__5ed539c67601c57bd107b5d8b04526be(actual_config_path: str, expected: Dict[str, str]) -> float:
    """
    Checks if VLC's video snapshot format is set to the expected value.
    The snapshot-format setting controls the image format for screenshots (png, jpg, etc.).
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_snapshot_format = expected['expected_snapshot_format']
    try:
        snapshot_format = 'png'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-format=' in line:
                snapshot_format = line.split('=')[-1].strip()
                break
        if snapshot_format == expected_snapshot_format:
            return 1.0
        else:
            logger.warning(f'Snapshot format mismatch - Expected: {expected_snapshot_format}, Found: {snapshot_format}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_video_on_top__9ca0ff77c9570a9a6e2f317bec352471(actual_config_path, rule):
    """
    Checks if VLC's video-on-top setting is enabled/disabled as expected.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_video_on_top = rule['expected_video_on_top']
    if isinstance(expected_video_on_top, int):
        expected_video_on_top = str(expected_video_on_top)
    try:
        video_on_top = '0'
        for line in config_file.split('\n'):
            if 'video-on-top=' in line:
                video_on_top = line.split('=')[-1].strip()
        if video_on_top == expected_video_on_top:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_aout__a69f9ac550f1e14ddbeccea862a27f73(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's audio output module is set to the expected value.
    The aout setting controls which audio output module VLC uses (pulse, alsa, etc.).
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_aout = rule['expected_aout']
    try:
        aout = None
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'aout=' in line and (not line.startswith('#')):
                aout = line.split('=')[-1].strip()
                break
        if aout is None:
            logger.info('No aout setting found in config, using default')
            if expected_aout in [None, '', 'default']:
                return 1.0
            else:
                return 0.0
        if aout == expected_aout:
            return 1.0
        else:
            logger.warning(f'Audio output module mismatch - Expected: {expected_aout}, Found: {aout}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_loop_mode__78fd4227(actual_config_path, expected):
    """
    Checks if VLC's loop mode is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_loop = expected.get('loop_mode')
    if isinstance(expected_loop, int):
        expected_loop = str(expected_loop)
    try:
        loop_mode = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'loop=' in line:
                loop_mode = line.split('=')[-1].strip()
                break
        if loop_mode == expected_loop:
            return 1.0
        else:
            logger.warning(f'Loop mode mismatch - Expected: {expected_loop}, Found: {loop_mode}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_continue__a10b6972(actual_config_path, expected):
    """
    Check if VLC is configured to continue playback (resume from last position).
    This checks the qt-continue setting in VLC config.

    Args:
        actual_config_path: Path to vlcrc config file
        expected: Expected value (0=ask, 1=always continue, 2=never continue)

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_continue = expected.get('expected_qt_continue')
    if isinstance(expected_qt_continue, int):
        expected_qt_continue = str(expected_qt_continue)
    try:
        qt_continue = '0'
        for line in config_file.split('\n'):
            if 'qt-continue=' in line:
                qt_continue = line.split('=')[-1].strip()
        if qt_continue == expected_qt_continue:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_aspect_ratio__d1470170(result, expected, **options):
    """
    Compare VLC aspect ratio against expected value.

    Args:
        result: Current aspect ratio from getter (str or None)
        expected: Expected aspect ratio (str, e.g., '16:9', '4:3')
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC aspect ratio result is None')
        return 0.0
    result_normalized = result.strip().lower()
    expected_normalized = expected.strip().lower()
    if result_normalized == expected_normalized:
        return 1.0
    else:
        logger.info(f'Aspect ratio mismatch - Expected: {expected}, Got: {result}')
        return 0.0

def check_audio_file_created__4071abfac3d4e1d773938a9f2b9279b5(result, expected, **options):
    """
    Check if audio file was created with correct properties.
    Expects: dict with 'filename', 'extension', 'min_size'
    Result: dict with 'exists', 'size', 'extension', 'filename'
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
        logger.info(f"File exists: {result.get('filename', 'unknown')}")
    else:
        logger.warning('File does not exist')
        return score
    expected_filename = expected.get('filename', '')
    if result.get('filename', '') == expected_filename:
        score += 0.3
        logger.info(f'Filename matches: {expected_filename}')
    else:
        logger.warning(f"Filename mismatch - Expected: {expected_filename}, Got: {result.get('filename', '')}")
    expected_ext = expected.get('extension', '')
    if result.get('extension', '') == expected_ext:
        score += 0.2
        logger.info(f'Extension matches: {expected_ext}')
    else:
        logger.warning(f"Extension mismatch - Expected: {expected_ext}, Got: {result.get('extension', '')}")
    min_size = expected.get('min_size', 1000)
    if result.get('size', 0) >= min_size:
        score += 0.1
        logger.info(f"File size OK: {result.get('size', 0)} bytes (min: {min_size})")
    else:
        logger.warning(f"File size too small: {result.get('size', 0)} bytes (min: {min_size})")
    return score

def check_vlc_playback_speed__c3b737f4(result, expected, **options):
    """
    Check if VLC playback speed matches the expected speed.

    Args:
        result: Current playback speed (float)
        expected: Expected configuration dict with 'speed' key
        **options: Additional options

    Returns:
        float: 1.0 if speed matches expected (within tolerance), 0.0 otherwise
    """
    expected_speed = expected.get('speed', 0.5)
    if abs(result - expected_speed) < 0.01:
        return 1.0
    else:
        return 0.0

def check_vlc_video_transformed__636928c8(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly transformed (rotated/flipped).

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename', 'source_filename', 'rotation_direction', and 'rotation_angle_degrees' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was transformed correctly,
               0.5 if transformation detected but direction cannot be verified,
               0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    expected_rotation_direction = expected.get('rotation_direction', 'counterclockwise')
    expected_rotation_angle = expected.get('rotation_angle_degrees', 90)
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    transformation_detected = False
    rotation_verified = False
    if output_filesize is not None and source_filesize is not None:
        if output_filesize != source_filesize:
            logger.info(f'File size changed: {source_filesize} -> {output_filesize}')
            transformation_detected = True
        else:
            logger.warning('File sizes identical - likely renamed/copied, not transformed')
    output_rotation = output_info.get('rotate', '0')
    source_rotation = source_info.get('rotate', '0') if source_info else '0'
    if expected_rotation_direction == 'counterclockwise' and expected_rotation_angle == 90:
        expected_output_rotation_values = [270, -90]
    elif expected_rotation_direction == 'clockwise' and expected_rotation_angle == 90:
        expected_output_rotation_values = [90]
    else:
        expected_output_rotation_values = []
        logger.warning(f'Unsupported rotation: {expected_rotation_direction} {expected_rotation_angle}°')
    try:
        output_rot_value = int(output_rotation) if output_rotation else 0
        source_rot_value = int(source_rotation) if source_rotation else 0
        if output_rot_value in expected_output_rotation_values:
            rotation_verified = True
            transformation_detected = True
            logger.info(f'Rotation verified: {source_rotation} -> {output_rotation} (expected {expected_rotation_direction} {expected_rotation_angle}°)')
        else:
            rotation_delta = (output_rot_value - source_rot_value) % 360
            if rotation_delta in expected_output_rotation_values:
                rotation_verified = True
                transformation_detected = True
                logger.info(f'Rotation verified via delta: {source_rotation} -> {output_rotation} (delta={rotation_delta})')
            else:
                logger.warning(f'Rotation metadata does not match expected {expected_rotation_direction} {expected_rotation_angle}°: {source_rotation} -> {output_rotation}')
    except (ValueError, TypeError):
        logger.warning(f'Could not parse rotation values: source={source_rotation}, output={output_rotation}')
    dimensions_swapped = False
    if source_info:
        output_width = output_info.get('width', '')
        output_height = output_info.get('height', '')
        source_width = source_info.get('width', '')
        source_height = source_info.get('height', '')
        if output_width == source_height and output_height == source_width and (output_width != output_height):
            logger.info(f'Dimensions swapped: {source_width}x{source_height} -> {output_width}x{output_height}')
            dimensions_swapped = True
            transformation_detected = True
    if rotation_verified and transformation_detected:
        logger.info('Video transformation fully verified with rotation metadata')
        return 1.0
    elif dimensions_swapped and transformation_detected and (not rotation_verified):
        logger.warning('Dimensions swapped but rotation direction cannot be verified from metadata - returning partial credit')
        return 0.5
    elif transformation_detected and (not dimensions_swapped) and (not rotation_verified):
        logger.warning('File was transformed but rotation cannot be verified')
        return 0.0
    else:
        logger.warning('No clear evidence of video transformation')
        return 0.0

def check_vlc_subtitle_autoload__989833e1(actual_config_path, expected):
    """
    Checks if VLC's subtitle autoload is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_autoload = expected.get('subtitle_autoload')
    if isinstance(expected_autoload, int):
        expected_autoload = str(expected_autoload)
    try:
        subtitle_autoload = '1'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'sub-autodetect-file=' in line:
                subtitle_autoload = line.split('=')[-1].strip()
                break
        if subtitle_autoload == expected_autoload:
            return 1.0
        else:
            logger.warning(f'Subtitle autoload mismatch - Expected: {expected_autoload}, Found: {subtitle_autoload}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_subtitle_setting__71783eb7(actual_config_path, expected):
    """Check VLC subtitle autodetect setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_sub_autodetect = expected.get('expected_sub_autodetect')
    if isinstance(expected_sub_autodetect, int):
        expected_sub_autodetect = str(expected_sub_autodetect)
    try:
        sub_autodetect = '1'
        for line in config_file.split('\n'):
            if 'sub-autodetect-file=' in line and (not line.strip().startswith('#')):
                sub_autodetect = line.split('=')[-1].strip()
        if sub_autodetect == expected_sub_autodetect:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_bookmarks_count__3288b82e(result, expected, **options):
    """
    Check if VLC has bookmarks created at the expected position.

    Args:
        result: Dictionary with 'count' and 'bookmarks' list
        expected: Expected configuration with 'min_bookmarks', 'expected_position', 'position_tolerance'
        **options: Additional options

    Returns:
        float: 1.0 if bookmarks exist at the expected position, 0.0 otherwise
    """
    min_bookmarks = expected.get('min_bookmarks', 1)
    expected_position = expected.get('expected_position')
    position_tolerance = expected.get('position_tolerance', 2.0)
    if isinstance(result, dict):
        bookmark_count = result.get('count', 0)
        bookmarks = result.get('bookmarks', [])
    else:
        bookmark_count = result
        bookmarks = []
    if bookmark_count < min_bookmarks:
        return 0.0
    if expected_position is not None and bookmarks:
        for bookmark in bookmarks:
            position = bookmark.get('position')
            if position is not None:
                try:
                    position_float = float(position)
                    if abs(position_float - expected_position) <= position_tolerance:
                        return 1.0
                except (ValueError, TypeError):
                    position_str = str(position).lower().strip()
                    if position_str.endswith('s'):
                        position_str = position_str[:-1]
                    if ':' in position_str:
                        try:
                            parts = position_str.split(':')
                            if len(parts) == 2:
                                minutes = float(parts[0])
                                seconds = float(parts[1])
                                position_float = minutes * 60 + seconds
                                if abs(position_float - expected_position) <= position_tolerance:
                                    return 1.0
                        except (ValueError, TypeError):
                            continue
                    else:
                        try:
                            position_float = float(position_str)
                            if abs(position_float - expected_position) <= position_tolerance:
                                return 1.0
                        except (ValueError, TypeError):
                            continue
        return 0.0
    return 1.0 if bookmark_count >= min_bookmarks else 0.0

def check_vlc_video_transformed__2331e840(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly transformed (rotated 90 degrees clockwise).

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was transformed correctly, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    expected_source_filename = expected.get('source_filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    remaining_after_output = probe_parts[1] if len(probe_parts) > 1 else ''
    source_parts = remaining_after_output.split('---SOURCE_LS---')
    source_video_info = source_parts[0].strip() if len(source_parts) > 0 else ''
    source_ls_output = source_parts[1].strip() if len(source_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    if not source_info:
        logger.warning('Could not read source video metadata - cannot verify transformation')
        return 0.0
    if expected_source_filename:
        if source_ls_output:
            if 'No such file or directory' in source_ls_output or 'cannot access' in source_ls_output:
                logger.warning(f'Source file does not exist: {expected_source_filename}')
                return 0.0
            if expected_source_filename not in source_ls_output:
                logger.warning(f"Source file '{expected_source_filename}' not found in source ls output")
                return 0.0
            logger.info(f'Source file validated: {expected_source_filename}')
        else:
            logger.info(f'Verifying against expected source: {expected_source_filename} (no ls output available)')
    transformation_detected = False
    rotation_correct = False
    if output_filesize is not None and source_filesize is not None:
        if output_filesize != source_filesize:
            logger.info(f'File size changed: {source_filesize} -> {output_filesize}')
            transformation_detected = True
        else:
            logger.warning('File sizes identical - likely renamed/copied, not transformed')
            return 0.0
    output_width = output_info.get('width', '')
    output_height = output_info.get('height', '')
    source_width = source_info.get('width', '')
    source_height = source_info.get('height', '')
    output_rotation = output_info.get('rotate', '0')
    source_rotation = source_info.get('rotate', '0')
    dimensions_swapped = output_width == source_height and output_height == source_width and (output_width != output_height)
    if dimensions_swapped:
        logger.info(f'Dimensions swapped correctly: {source_width}x{source_height} -> {output_width}x{output_height}')
        try:
            output_rot_val = int(output_rotation)
            source_rot_val = int(source_rotation)
            if output_rot_val == 90:
                rotation_correct = True
                logger.info(f'Rotation angle verified: 90 degrees clockwise (rotate={output_rot_val})')
            elif output_rot_val == 0 and dimensions_swapped:
                rotation_correct = True
                logger.info(f'Rotation verified: dimensions swapped, rotate=0 (physical rotation)')
            elif output_rot_val == 270 or output_rot_val == -90:
                logger.warning(f'Rotation is 270 degrees counterclockwise (rotate={output_rot_val}), not 90 clockwise')
                rotation_correct = False
            elif output_rot_val == 180 or output_rot_val == -180:
                logger.warning(f'Rotation is 180 degrees (rotate={output_rot_val}), not 90 clockwise')
                rotation_correct = False
            else:
                logger.warning(f'Unexpected rotation metadata value: {output_rot_val}. Cannot verify rotation direction without valid metadata.')
                rotation_correct = False
        except (ValueError, TypeError):
            logger.warning(f"Rotation metadata missing or invalid (output: '{output_rotation}', source: '{source_rotation}'). Cannot distinguish 90° clockwise from 270° counterclockwise.")
            rotation_correct = False
    else:
        logger.warning('Dimensions not swapped - video may not be rotated 90 degrees')
        rotation_correct = False
    if output_rotation != source_rotation:
        logger.info(f'Rotation metadata changed: {source_rotation} -> {output_rotation}')
        transformation_detected = True
    if transformation_detected and rotation_correct:
        logger.info('Video transformation verified: 90 degree clockwise rotation detected')
        return 1.0
    elif transformation_detected and (not rotation_correct):
        logger.warning('Video was transformed but rotation direction is incorrect')
        return 0.0
    else:
        logger.warning('No clear evidence of video transformation')
        return 0.0

def check_vlc_subtitle_disabled__65753b1a(result, expected, **options):
    """
    Check if VLC subtitles are disabled.

    Args:
        result: Current subtitle track ID from getter (int or None)
        expected: Expected disabled status - True means should be disabled (track < 0)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC subtitle track result is None')
        return 0.0
    is_disabled = result < 0
    if is_disabled == expected:
        return 1.0
    else:
        logger.info(f'Subtitle status mismatch - Expected disabled: {expected}, Track ID: {result}')
        return 0.0

def check_vlc_repeat_mode__8e82e8fa6a0ee6159e766778349645ac(actual_config_path, rule):
    """
    Checks if VLC's repeat mode is set to the expected value.
    Repeat mode can be: 0 (no repeat), 1 (repeat current), 2 (repeat all)
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_repeat = rule['expected_repeat']
    if isinstance(expected_repeat, int):
        expected_repeat = str(expected_repeat)
    try:
        qt_continue = '0'
        for line in config_file.split('\n'):
            if 'qt-continue=' in line:
                qt_continue = line.split('=')[-1].strip()
        if qt_continue == expected_repeat:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_qt_start_minimized__b23ac353(actual_config_path, rule):
    """Check if VLC start minimized setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_start_minimized = rule['expected_qt_start_minimized']
    if isinstance(expected_qt_start_minimized, int):
        expected_qt_start_minimized = str(expected_qt_start_minimized)
    try:
        qt_start_minimized = '0'
        for line in config_file.split('\n'):
            if 'qt-start-minimized=' in line:
                qt_start_minimized = line.split('=')[-1].strip()
        if qt_start_minimized == expected_qt_start_minimized:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_video_output_dir__8b4751b6(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's video output directory is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_video_dir = rule['video_output_dir']
    try:
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'sout-file-append' in line or 'sout-standard-path' in line:
                current_path = line.split('=')[-1].strip()
                if current_path == expected_video_dir:
                    return 1.0
                else:
                    return 0.0
        return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_audio_info__c6746b3e(result, expected, **options):
    """
    Check if the audio info file contains the expected content.

    Args:
        result: Path to the audio_info.txt file
        expected: Dict with "expected_content" key
        **options: Additional options

    Returns:
        float: 1.0 if correct, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        expected_content = expected.get('expected_content', '')
        if content == expected_content:
            return 1.0
        else:
            logger.info(f'Content mismatch. Expected: {expected_content}, Got: {content}')
            return 0.0
    except Exception as e:
        logger.error(f'Error reading file: {e}')
        return 0.0

def check_mp3_audio_extraction__5e12822f6a68285e2467088dd7d25598(result, expected, **options):
    """
    Comprehensive check for MP3 audio extraction from video.
    Verifies file format, audio content, quality, and source-target relationship.

    Args:
        result: Dict from get_mp3_audio_info__5e12822f6a68285e2467088dd7d25598
        expected: Dict with expected values:
            - min_size_mb: minimum file size in MB
            - max_size_mb: maximum file size in MB
            - check_duration_match: whether to verify duration matches source (default True)
            - require_mp3_format: whether to strictly verify MP3 format (default True)
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on how many checks pass
    """
    min_size_mb = expected.get('min_size_mb', 1.0)
    max_size_mb = expected.get('max_size_mb', 10.0)
    check_duration_match = expected.get('check_duration_match', True)
    require_mp3_format = expected.get('require_mp3_format', True)
    checks_passed = 0
    total_checks = 0
    issues = []
    total_checks += 1
    if result.get('exists', False):
        checks_passed += 1
        logger.info('✓ File exists')
    else:
        issues.append('File does not exist')
        logger.error('✗ File does not exist')
        return 0.0
    if require_mp3_format:
        total_checks += 1
        if result.get('is_mp3', False):
            checks_passed += 1
            logger.info(f"✓ File is MP3 format: {result.get('file_type', '')}")
        else:
            issues.append(f"File is not MP3 format. Detected type: {result.get('file_type', 'unknown')}")
            logger.error(f"✗ File is not MP3 format: {result.get('file_type', '')}")
    total_checks += 1
    if result.get('has_audio', False):
        checks_passed += 1
        logger.info(f"✓ File contains valid audio stream (codec: {result.get('codec', 'unknown')})")
    else:
        issues.append('File does not contain valid audio stream')
        logger.error('✗ File does not contain valid audio stream')
    total_checks += 1
    codec = result.get('codec', '').lower()
    if 'mp3' in codec or codec in ['mp2', 'mp1', 'mpeg1audio', 'mpeg2audio']:
        checks_passed += 1
        logger.info(f"✓ Audio codec is MP3-compatible: {result.get('codec', '')}")
    else:
        issues.append(f"Audio codec is not MP3: {result.get('codec', 'unknown')}")
        logger.error(f"✗ Audio codec is not MP3: {result.get('codec', '')}")
    total_checks += 1
    size_mb = result.get('size_mb', 0.0)
    if min_size_mb <= size_mb <= max_size_mb:
        checks_passed += 1
        logger.info(f'✓ File size {size_mb} MB is within range [{min_size_mb}, {max_size_mb}] MB')
    else:
        issues.append(f'File size {size_mb} MB is outside range [{min_size_mb}, {max_size_mb}] MB')
        logger.error(f'✗ File size {size_mb} MB is outside range [{min_size_mb}, {max_size_mb}] MB')
    total_checks += 1
    duration = result.get('duration', 0.0)
    if duration > 0:
        checks_passed += 1
        logger.info(f'✓ Audio has valid duration: {duration:.2f} seconds')
    else:
        issues.append('Audio duration is zero or invalid')
        logger.error('✗ Audio duration is zero or invalid')
    if check_duration_match and result.get('source_duration', 0.0) > 0:
        total_checks += 1
        if result.get('duration_match', False):
            checks_passed += 1
            logger.info(f"✓ Audio duration ({result.get('duration', 0):.2f}s) matches source video ({result.get('source_duration', 0):.2f}s)")
        else:
            issues.append(f"Audio duration ({result.get('duration', 0):.2f}s) does not match source video ({result.get('source_duration', 0):.2f}s)")
            logger.error(f"✗ Duration mismatch: MP3={result.get('duration', 0):.2f}s, Source={result.get('source_duration', 0):.2f}s")
    total_checks += 1
    sample_rate = result.get('sample_rate', 0)
    if sample_rate >= 22050:
        checks_passed += 1
        logger.info(f'✓ Audio has reasonable sample rate: {sample_rate} Hz')
    else:
        issues.append(f'Audio sample rate too low or invalid: {sample_rate} Hz')
        logger.error(f'✗ Audio sample rate too low: {sample_rate} Hz')
    if total_checks == 0:
        score = 0.0
    else:
        score = checks_passed / total_checks
    logger.info(f'MP3 audio extraction check: {checks_passed}/{total_checks} checks passed (score: {score:.2f})')
    if issues:
        logger.warning(f"Issues found: {', '.join(issues)}")
    if score == 1.0:
        return 1.0
    elif score >= 0.75:
        return score
    else:
        return 0.0

def check_vlc_qt_autoload_extensions__8d3419de(actual_config_path, expected):
    """Check VLC autoload extensions setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_autoload = expected.get('expected_autoload_extensions')
    if isinstance(expected_autoload, int):
        expected_autoload = str(expected_autoload)
    try:
        autoload_extensions = '1'
        for line in config_file.split('\n'):
            if 'qt-autoload-extensions=' in line and (not line.strip().startswith('#')):
                autoload_extensions = line.split('=')[-1].strip()
        if autoload_extensions == expected_autoload:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_time_advanced__2bf4fa278b8d46abde1b1e5c8949b96a(actual_status_path: str, expected: Dict) -> float:
    """
    Checks if VLC playback time has advanced beyond a minimum threshold.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Dict with min_time (minimum time in seconds the video should be at)

    Returns:
        1.0 if current time >= min_time, 0.0 otherwise
    """
    if not actual_status_path:
        logger.error('No VLC status path provided')
        return 0.0
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        time_element = tree.find('time')
        if time_element is not None:
            actual_time = int(time_element.text)
        else:
            actual_time = 0
        min_time = int(expected.get('min_time', 30))
        logger.info(f'VLC Playback Time: {actual_time}s, Minimum Expected: {min_time}s')
        if actual_time >= min_time:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC playback time: {e}')
        return 0.0

def check_video_duration_saved__da7e9a1f(result_state: Optional[str], expected_state: Dict[str, bool], **options) -> float:
    """
    Check if the video duration was correctly saved to duration.txt file.

    Scoring breakdown:
    - 0.5 for file exists with content
    - 0.5 for correct duration format (time information present)

    Expected video duration is approximately 122 seconds or 2:02 based on metadata.

    Args:
        result_state: Path to the cached duration.txt file, or None if file doesn't exist
        expected_state: Dict with rules, e.g., {"has_duration": true}
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result_state is None:
        logger.warning('Duration file does not exist or could not be retrieved')
        return 0.0
    try:
        with open(result_state, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read().strip()
        if not content:
            logger.warning('Duration file is empty')
            return 0.0
        score += 0.5
        logger.info(f'Duration file exists with content: {content}')
        time_patterns = ['\\d+:\\d+', '\\b\\d+\\s*(s|sec|second|seconds)\\b', '^\\d+$']
        has_time_info = False
        for pattern in time_patterns:
            if re.search(pattern, content, re.IGNORECASE):
                has_time_info = True
                logger.info(f"Found time information matching pattern '{pattern}': {content}")
                break
        if has_time_info:
            score += 0.5
        else:
            logger.warning(f'File content does not appear to contain duration/time information: {content}')
        return score
    except Exception as e:
        logger.error(f'Error reading duration file: {e}')
        return score

def check_audio_conversion__70ff6e4ac087c63cd4c29ec50188d44d(result, expected, **options):
    """
    Verify audio file conversion was successful.
    Checks file existence, naming, format, and content validity.
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.3
        logger.info(f"Audio file exists: {result.get('basename', 'unknown')}")
    else:
        logger.warning('Audio file does not exist')
        return score
    expected_basename = expected.get('basename', '')
    if result.get('basename', '') == expected_basename:
        score += 0.25
        logger.info(f'Basename matches: {expected_basename}')
    else:
        logger.warning(f"Basename mismatch - Expected: {expected_basename}, Got: {result.get('basename', '')}")
    if result.get('is_audio', False):
        score += 0.25
        logger.info(f"File is audio format: {result.get('mime_type', 'unknown')}")
    else:
        logger.warning(f"File is not audio format: {result.get('mime_type', 'unknown')}")
    min_size = expected.get('min_size', 50000)
    if result.get('size', 0) >= min_size:
        score += 0.2
        logger.info(f"File size sufficient: {result.get('size', 0)} bytes")
    else:
        logger.warning(f"File size too small: {result.get('size', 0)} bytes (min: {min_size})")
    return score

def check_vlc_video_on_top__81895b32(actual_config_path, rule):
    """
    Checks if VLC's video-on-top setting is configured correctly.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_video_on_top = rule['expected_video_on_top']
    if isinstance(expected_video_on_top, int):
        expected_video_on_top = str(expected_video_on_top)
    try:
        video_on_top = '0'
        for line in config_file.split('\n'):
            if 'video-on-top=' in line:
                video_on_top = line.split('=')[-1].strip()
        if video_on_top == expected_video_on_top:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_qt_privacy_ask__f900ec1c(actual_config_path, expected):
    """Check VLC privacy/network interaction setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_privacy_ask = expected.get('expected_privacy_ask')
    if isinstance(expected_privacy_ask, int):
        expected_privacy_ask = str(expected_privacy_ask)
    try:
        privacy_ask = '1'
        for line in config_file.split('\n'):
            if 'qt-privacy-ask=' in line and (not line.strip().startswith('#')):
                privacy_ask = line.split('=')[-1].strip()
        if privacy_ask == expected_privacy_ask:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_updates_notif__b84bb567(actual_config_path, rule):
    """Check if VLC update notification setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_updates_notif = rule['expected_qt_updates_notif']
    if isinstance(expected_qt_updates_notif, int):
        expected_qt_updates_notif = str(expected_qt_updates_notif)
    try:
        qt_updates_notif = '1'
        for line in config_file.split('\n'):
            if 'qt-updates-notif=' in line:
                qt_updates_notif = line.split('=')[-1].strip()
        if qt_updates_notif == expected_qt_updates_notif:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_continue_playback_enabled(actual_config_path: str, expected: Dict[str, str]) -> float:
    """
    Checks if VLC's continue playback feature is enabled in the configuration file.

    VLC uses the following settings in vlcrc for continue playback:
    - qt-continue=1: Continue playback (Qt interface on Linux/Windows)
    - macosx-continue-playback=1: Continue playback always (macOS)

    The default behavior is to ask (qt-continue=0 or not set).

    Args:
        actual_config_path: Path to the VLC config file
        expected: Dictionary with 'expected_continue_playback' key (0, 1, or 2)
              For qt-continue: 0=Ask (default), 1=Always continue
              For macosx-continue-playback: 0=Ask, 1=Always, 2=Never

    Returns:
        1.0 if continue playback is enabled as expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_value = expected.get('expected_continue_playback', '1')
    if isinstance(expected_value, int):
        expected_value = str(expected_value)
    try:
        qt_continue = None
        macosx_continue = None
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if line.startswith('qt-continue='):
                qt_continue = line.split('=')[-1].strip()
            elif line.startswith('macosx-continue-playback='):
                macosx_continue = line.split('=')[-1].strip()
        if qt_continue is not None:
            actual_value = qt_continue
        elif macosx_continue is not None:
            actual_value = macosx_continue
        else:
            actual_value = '0'
        logger.info(f'VLC continue playback status - qt-continue={qt_continue}, macosx-continue-playback={macosx_continue}, actual={actual_value}, expected={expected_value}')
        if actual_value == expected_value:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_continue_playback__90a64484(actual_config_path, rule):
    """Check if VLC 'Continue playback' setting is configured correctly."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_continue_playback = rule['expected_qt_continue_playback']
    if isinstance(expected_qt_continue_playback, int):
        expected_qt_continue_playback = str(expected_qt_continue_playback)
    try:
        qt_continue_playback = '1'
        for line in config_file.split('\n'):
            if 'qt-continue=' in line:
                qt_continue_playback = line.split('=')[-1].strip()
        if qt_continue_playback == expected_qt_continue_playback:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_vlc_key_fullscreen__a57dec40610279f297f0f0fecc93f9ff(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's fullscreen hotkey is set to the expected value.
    The key-toggle-fullscreen setting controls the keyboard shortcut for fullscreen toggle.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_key_fullscreen = rule['expected_key_fullscreen']
    try:
        key_fullscreen = 'f'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'key-toggle-fullscreen=' in line:
                key_fullscreen = line.split('=')[-1].strip()
                break
        if key_fullscreen == expected_key_fullscreen:
            return 1.0
        else:
            logger.warning(f'Fullscreen key mismatch - Expected: {expected_key_fullscreen}, Found: {key_fullscreen}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_video_on_top__901155ec(actual_config_path, expected):
    """
    Checks if VLC's video-on-top setting is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_on_top = expected.get('video_on_top')
    if isinstance(expected_on_top, int):
        expected_on_top = str(expected_on_top)
    try:
        video_on_top = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'video-on-top=' in line:
                video_on_top = line.split('=')[-1].strip()
                break
        if video_on_top == expected_on_top:
            return 1.0
        else:
            logger.warning(f'Video-on-top mismatch - Expected: {expected_on_top}, Found: {video_on_top}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_is_video_file__4939fb90(result, expected, **options):
    """Check if file exists and is a video file.

    Args:
        result: File properties dict from getter
        expected: Expected rules dict
        **options: Additional comparison options

    Returns:
        float: 1.0 if file exists and is video, 0.0 otherwise
    """
    if not result.get('exists', False):
        logger.warning('File does not exist')
        return 0.0
    file_type = result.get('file_type', '').lower()
    video_indicators = ['mp4', 'video', 'iso media', 'mpeg']
    is_video = any((indicator in file_type for indicator in video_indicators))
    if is_video:
        return 1.0
    else:
        logger.info(f"File type '{file_type}' does not appear to be a video file")
        return 0.0

def check_snapshot_created__4e252238(result_state: Optional[Dict], expected_state: Dict, **options) -> float:
    """
    Check if a snapshot file was created correctly from VLC.

    Args:
        result_state: Dict with file properties from getter:
            {
                'exists': bool,
                'is_image': bool,
                'size': int,
                'path': str
            }
        expected_state: Dict with expected properties:
            {
                'exists': bool (expected existence state),
                'is_image': bool (expected image format state)
            }
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
            - 0.4 for file existence matching expected
            - 0.4 for being an image file matching expected
            - 0.2 for valid file size (> 1KB and < 100MB)
    """
    if result_state is None:
        logger.error('Result state is None')
        return 0.0
    score = 0.0
    expected_exists = expected_state.get('exists', True)
    expected_is_image = expected_state.get('is_image', True)
    result_exists = result_state.get('exists', False)
    if result_exists == expected_exists:
        score += 0.4
        logger.info(f'Snapshot file existence matches expected ({expected_exists}): +0.4')
    else:
        logger.info(f'Snapshot file existence ({result_exists}) does not match expected ({expected_exists}): +0.0')
    result_is_image = result_state.get('is_image', False)
    if result_is_image == expected_is_image:
        score += 0.4
        logger.info(f'File image format matches expected ({expected_is_image}): +0.4')
    else:
        logger.info(f'File image format ({result_is_image}) does not match expected ({expected_is_image}): +0.0')
    if result_exists and expected_exists:
        file_size = result_state.get('size', 0)
        if file_size > 1024 and file_size < 100 * 1024 * 1024:
            score += 0.2
            logger.info(f'File size is valid ({file_size} bytes): +0.2')
        else:
            logger.info(f'File size is invalid ({file_size} bytes): +0.0')
    elif result_exists == expected_exists:
        score += 0.2
        logger.info("File size check skipped (file doesn't exist as expected): +0.2")
    logger.info(f'Total score: {score}')
    return score

def check_audio_format__5d993657(result, expected, **options):
    """Check if audio file has the expected format, filename, and exists.

    Args:
        result: dict with 'file_type', 'is_mp3', 'file_exists', and 'filename_correct' keys
        expected: dict with 'is_mp3' boolean

    Returns:
        float: 1.0 if all checks pass, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.warning(f'Result is not a dict: {type(result)}')
        return 0.0
    file_exists = result.get('file_exists', False)
    if not file_exists:
        logger.warning(f"File does not exist at path: {result.get('path', 'unknown')}")
        return 0.0
    filename_correct = result.get('filename_correct', False)
    if not filename_correct:
        logger.warning(f"Filename incorrect: expected '{result.get('expected_filename', 'unknown')}', got '{result.get('actual_filename', 'unknown')}'")
        return 0.0
    is_mp3 = result.get('is_mp3', False)
    expected_is_mp3 = expected.get('is_mp3', True)
    if is_mp3 == expected_is_mp3:
        logger.info(f"All checks passed: file exists, filename correct, MP3={is_mp3}, file_type={result.get('file_type', 'unknown')}")
        return 1.0
    else:
        logger.info(f'Audio format incorrect: expected MP3={expected_is_mp3}, got MP3={is_mp3}')
        return 0.0

def check_snapshot_size__313dd5e1(result, expected, **options):
    """Check if snapshot file size is within expected range.

    Args:
        result: File size in bytes from getter
        expected: Size constraints from rules
        **options: Additional options

    Returns:
        1.0 if size is valid, 0.0 otherwise
    """
    min_size = expected.get('min_size', 10000)
    max_size = expected.get('max_size', 10000000)
    if min_size <= result <= max_size:
        return 1.0
    else:
        return 0.0

def check_audio_file_format__6475bf6e7aa0ef4599a6c11ec95f5406(result, expected, **options):
    """
    Check if audio file exists and has expected format.

    Args:
        result: dict from getter with 'exists', 'format', 'size' keys
        expected: dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    score += 0.4
    expected_format = expected.get('format')
    actual_format = result.get('format')
    if expected_format and actual_format == expected_format:
        score += 0.3
        logger.info(f'Format matches: {actual_format}')
    elif expected_format:
        logger.warning(f'Format mismatch - expected: {expected_format}, got: {actual_format}')
    min_size = expected.get('min_size', 0)
    actual_size = result.get('size', 0)
    if actual_size >= min_size:
        score += 0.3
        logger.info(f'Size is reasonable: {actual_size} >= {min_size}')
    else:
        logger.warning(f'Size too small: {actual_size} < {min_size}')
    return score

def check_vlc_start_paused__443cb77a(actual_config_path, expected):
    """
    Checks if VLC's start paused setting is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected configuration (from rules dict)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_start_paused = expected['expected_start_paused']
    if isinstance(expected_start_paused, int):
        expected_start_paused = str(expected_start_paused)
    try:
        start_paused = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'start-paused=' in line:
                start_paused = line.split('=')[-1].strip()
        if start_paused == expected_start_paused:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_qt_notification__efa02ee6(actual_config_path, expected):
    """Check VLC popup notification setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_notification = expected.get('expected_notification')
    if isinstance(expected_notification, int):
        expected_notification = str(expected_notification)
    try:
        notification = '1'
        for line in config_file.split('\n'):
            if 'qt-notification=' in line and (not line.strip().startswith('#')):
                notification = line.split('=')[-1].strip()
        if notification == expected_notification:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_video_transformed__93a4e125(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly horizontally flipped.

    Verifies horizontal flip by checking:
    1. File exists with correct name
    2. File was re-encoded (different size)
    3. Dimensions remain the same (horizontal flip preserves width x height)
    4. CRITICAL: Left/right edge pixels are swapped (pixel-level verification)

    Args:
        result: Combined output from ls, stat, ffprobe, and edge pixel hash commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was horizontally flipped, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.0
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.0
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.0
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    remaining_after_source = probe_parts[1] if len(probe_parts) > 1 else ''
    edge_parts = remaining_after_source.split('---EDGE_DATA---')
    source_video_info = edge_parts[0].strip() if len(edge_parts) > 0 else ''
    src_left_f0 = ''
    out_right_f0 = ''
    src_right_f0 = ''
    out_left_f0 = ''
    src_left_f30 = ''
    out_right_f30 = ''
    src_right_f30 = ''
    out_left_f30 = ''
    if len(edge_parts) > 1:
        remaining = edge_parts[1]
        parts_list = remaining.split('---SRC_LEFT_F0---')
        if len(parts_list) > 1:
            remaining = parts_list[1]
            parts_list = remaining.split('---OUT_RIGHT_F0---')
            src_left_f0 = parts_list[0].strip() if len(parts_list) > 0 else ''
            if len(parts_list) > 1:
                remaining = parts_list[1]
                parts_list = remaining.split('---SRC_RIGHT_F0---')
                out_right_f0 = parts_list[0].strip() if len(parts_list) > 0 else ''
                if len(parts_list) > 1:
                    remaining = parts_list[1]
                    parts_list = remaining.split('---OUT_LEFT_F0---')
                    src_right_f0 = parts_list[0].strip() if len(parts_list) > 0 else ''
                    if len(parts_list) > 1:
                        remaining = parts_list[1]
                        parts_list = remaining.split('---SRC_LEFT_F30---')
                        out_left_f0 = parts_list[0].strip() if len(parts_list) > 0 else ''
                        if len(parts_list) > 1:
                            remaining = parts_list[1]
                            parts_list = remaining.split('---OUT_RIGHT_F30---')
                            src_left_f30 = parts_list[0].strip() if len(parts_list) > 0 else ''
                            if len(parts_list) > 1:
                                remaining = parts_list[1]
                                parts_list = remaining.split('---SRC_RIGHT_F30---')
                                out_right_f30 = parts_list[0].strip() if len(parts_list) > 0 else ''
                                if len(parts_list) > 1:
                                    remaining = parts_list[1]
                                    parts_list = remaining.split('---OUT_LEFT_F30---')
                                    src_right_f30 = parts_list[0].strip() if len(parts_list) > 0 else ''
                                    if len(parts_list) > 1:
                                        out_left_f30 = parts_list[1].strip()
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    logger.info('Output is a valid video file')
    if source_info:
        output_width = output_info.get('width', '')
        output_height = output_info.get('height', '')
        source_width = source_info.get('width', '')
        source_height = source_info.get('height', '')
        if not (output_width and output_height and source_width and source_height):
            logger.warning('Could not determine video dimensions')
            return 0.0
        if output_width != source_width or output_height != source_height:
            logger.warning(f'Dimensions changed: {source_width}x{source_height} -> {output_width}x{output_height}')
            logger.warning('This indicates rotation or cropping, not horizontal flip')
            return 0.0
        logger.info(f'Dimensions preserved: {output_width}x{output_height} (consistent with horizontal flip)')

    def extract_hash(hash_str):
        if not hash_str or 'HASH_ERROR' in hash_str or 'ERROR' in hash_str:
            return None
        parts = hash_str.strip().split()
        return parts[0] if parts else None
    src_left_f0_hash = extract_hash(src_left_f0)
    out_right_f0_hash = extract_hash(out_right_f0)
    src_right_f0_hash = extract_hash(src_right_f0)
    out_left_f0_hash = extract_hash(out_left_f0)
    src_left_f30_hash = extract_hash(src_left_f30)
    out_right_f30_hash = extract_hash(out_right_f30)
    src_right_f30_hash = extract_hash(src_right_f30)
    out_left_f30_hash = extract_hash(out_left_f30)
    flip_verified_f0 = False
    if src_left_f0_hash and out_right_f0_hash and src_right_f0_hash and out_left_f0_hash:
        if src_left_f0_hash == out_right_f0_hash and src_right_f0_hash == out_left_f0_hash:
            logger.info('Frame 0: Horizontal flip VERIFIED - edge pixels are correctly swapped')
            flip_verified_f0 = True
        else:
            logger.warning(f'Frame 0: Edge pixels NOT swapped correctly')
            logger.warning(f'  Source left: {src_left_f0_hash}, Output right: {out_right_f0_hash} - Match: {src_left_f0_hash == out_right_f0_hash}')
            logger.warning(f'  Source right: {src_right_f0_hash}, Output left: {out_left_f0_hash} - Match: {src_right_f0_hash == out_left_f0_hash}')
    else:
        logger.warning('Frame 0: Could not extract edge pixel data')
    flip_verified_f30 = False
    if src_left_f30_hash and out_right_f30_hash and src_right_f30_hash and out_left_f30_hash:
        if src_left_f30_hash == out_right_f30_hash and src_right_f30_hash == out_left_f30_hash:
            logger.info('Frame 30: Horizontal flip VERIFIED - edge pixels are correctly swapped')
            flip_verified_f30 = True
        else:
            logger.warning(f'Frame 30: Edge pixels NOT swapped correctly')
            logger.warning(f'  Source left: {src_left_f30_hash}, Output right: {out_right_f30_hash} - Match: {src_left_f30_hash == out_right_f30_hash}')
            logger.warning(f'  Source right: {src_right_f30_hash}, Output left: {out_left_f30_hash} - Match: {src_right_f30_hash == out_left_f30_hash}')
    else:
        logger.warning('Frame 30: Could not extract edge pixel data')
    if flip_verified_f0 or flip_verified_f30:
        logger.info('HORIZONTAL FLIP VERIFIED: Edge pixel analysis confirms horizontal flip was applied')
        return 1.0
    else:
        logger.warning('HORIZONTAL FLIP NOT VERIFIED: Edge pixel analysis does not confirm horizontal flip')
        logger.warning('The video may have been transformed differently (brightness, color, etc.) but NOT horizontally flipped')
        return 0.0

def check_vlc_playlist_count__fc73aaff(result, expected, **options):
    """
    Check if VLC playlist has more than one item.

    Args:
        result: Number of items in playlist
        expected: Expected configuration
        **options: Additional options

    Returns:
        float: 1.0 if playlist has multiple items, 0.0 otherwise
    """
    min_items = expected.get('min_items', 2)
    if result >= min_items:
        return 1.0
    else:
        return 0.0

def check_vlc_video_transformed__fc9f20bb(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly transformed (rotated/flipped).

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was transformed, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    transformation_detected = False
    if output_filesize is not None and source_filesize is not None:
        if output_filesize != source_filesize:
            logger.info(f'File size changed: {source_filesize} -> {output_filesize}')
            transformation_detected = True
        else:
            logger.warning('File sizes identical - likely renamed/copied, not transformed')
    output_rotation = output_info.get('rotate', '0')
    source_rotation = source_info.get('rotate', '0') if source_info else '0'
    if output_rotation != source_rotation:
        logger.info(f'Rotation changed: {source_rotation} -> {output_rotation}')
        transformation_detected = True
    if source_info:
        output_width = output_info.get('width', '')
        output_height = output_info.get('height', '')
        source_width = source_info.get('width', '')
        source_height = source_info.get('height', '')
        if output_width == source_height and output_height == source_width and (output_width != output_height):
            logger.info(f'Dimensions swapped: {source_width}x{source_height} -> {output_width}x{output_height}')
            transformation_detected = True
    if transformation_detected:
        logger.info('Video transformation verified')
        return 1.0
    else:
        logger.warning('No clear evidence of video transformation')
        return 0.0

def check_snapshot_saved__0706c584(result, expected, **options):
    """Check if a snapshot file was successfully saved.

    Args:
        result: Boolean indicating if file exists (from getter)
        expected: Expected value (should be True)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    expected_value = expected.get('exists', True)
    logger.info(f'Snapshot file exists: {result}, expected: {expected_value}')
    if result == expected_value:
        return 1.0
    else:
        return 0.0

def check_audio_extracted__d0f84a157daee05fd0a57bb243a6ea5c(result, expected, **options):
    """
    Check if audio was properly extracted from video.

    Args:
        result: Dict with audio file status from getter
        expected: Dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.5
    if result.get('size', 0) > 10000:
        score += 0.25
    if result.get('duration', 0) > 1.0:
        score += 0.25
    return score

def check_snapshot_created__1ce093b98b48496fecfb155d31ed7704(result, expected, **options):
    """
    Check if snapshot was properly created from video.

    Args:
        result: Dict with snapshot status from getter
        expected: Dict with expected values
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    if result.get('is_image', False):
        score += 0.4
    if result.get('size', 0) > 1000:
        score += 0.2
    return score

def check_vlc_playback_rate__d5f08d8dc19ae133ba902b1660141393(actual_status_path: str, expected: Dict) -> float:
    """
    Checks if VLC playback rate/speed matches expected value.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Dict with expected_rate (e.g., 1.5 for 1.5x speed, 0.5 for half speed)

    Returns:
        1.0 if playback rate matches expected (within tolerance), 0.0 otherwise
    """
    if not actual_status_path:
        logger.error('No VLC status path provided')
        return 0.0
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        rate_element = tree.find('rate')
        if rate_element is not None:
            actual_rate = float(rate_element.text)
        else:
            actual_rate = 1.0
        expected_rate = float(expected.get('expected_rate', 1.0))
        tolerance = float(expected.get('tolerance', 0.01))
        logger.info(f'VLC Playback Rate: {actual_rate}, Expected: {expected_rate}')
        if abs(actual_rate - expected_rate) <= tolerance:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC playback rate: {e}')
        return 0.0

def check_vlc_random_mode__de28db3c(actual_config_path, expected):
    """
    Checks if VLC's random playback mode is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_random = expected.get('random_mode')
    if isinstance(expected_random, int):
        expected_random = str(expected_random)
    try:
        random_mode = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'random=' in line:
                random_mode = line.split('=')[-1].strip()
                break
        if random_mode == expected_random:
            return 1.0
        else:
            logger.warning(f'Random mode mismatch - Expected: {expected_random}, Found: {random_mode}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_snapshot_directory__008263c0(result, expected, **options):
    """Check if VLC snapshot directory matches expected path.

    Args:
        result: Current snapshot directory from getter
        expected: Expected directory from rules
        **options: Additional options

    Returns:
        1.0 if matches, 0.0 otherwise
    """
    expected_dir = expected.get('directory', '/home/user/Desktop')
    if result == expected_dir:
        return 1.0
    else:
        return 0.0

def check_qt_continue_playback__e7d03ebf(actual_config_path, expected, **options):
    """
    Checks if VLC's continue playback setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_continue_playback" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_continue_playback = expected['expected_qt_continue_playback']
    if isinstance(expected_qt_continue_playback, int):
        expected_qt_continue_playback = str(expected_qt_continue_playback)
    try:
        qt_continue_playback = '0'
        for line in config_file.split('\n'):
            if 'qt-continue=' in line:
                qt_continue_playback = line.split('=')[-1].strip()
        if qt_continue_playback == expected_qt_continue_playback:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_qt_fs_controller__dd2b0e06(actual_config_path, expected):
    """Check VLC fullscreen controller setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_fs_controller = expected.get('expected_fs_controller')
    if isinstance(expected_fs_controller, int):
        expected_fs_controller = str(expected_fs_controller)
    try:
        fs_controller = '1'
        for line in config_file.split('\n'):
            if 'qt-fs-controller=' in line and (not line.strip().startswith('#')):
                fs_controller = line.split('=')[-1].strip()
        if fs_controller == expected_fs_controller:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_rotation__ad793600129eb921cf29c68b343191af(result, expected, **options):
    """Check if video has been properly rotated.

    Args:
        result: Dict from getter with video metadata:
            {
                'exists': bool,
                'path': str,
                'rotation': int (0, 90, 180, 270),
                'width': int,
                'height': int,
                'file_size': int
            }
        expected: Dict with expected video properties:
            {
                'exists': bool,
                'rotation': int (0, 90, 180, 270),
                'min_size': int (minimum file size to ensure it's not just a copy)
            }
        **options: Additional options

    Returns:
        float: 1.0 if video meets all criteria, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Invalid result format: {result}')
        return 0.0
    if not result.get('exists', False):
        logger.warning(f"Video file does not exist: {result.get('path', 'unknown')}")
        return 0.0
    expected_exists = expected.get('exists', True)
    if not expected_exists:
        logger.info(f"Video file exists but should not: {result.get('path', 'unknown')}")
        return 0.0
    expected_rotation = expected.get('rotation')
    actual_rotation = result.get('rotation')
    if expected_rotation is not None:
        if actual_rotation is None:
            logger.warning(f"Could not determine rotation for {result.get('path', 'unknown')}")
            if result.get('file_size', 0) > expected.get('min_size', 0):
                logger.info(f'Video exists with reasonable size, but rotation unclear - giving partial credit')
                return 0.7
            return 0.0
        expected_rotation = int(expected_rotation) % 360
        actual_rotation = int(actual_rotation) % 360
        if actual_rotation != expected_rotation:
            logger.warning(f"Video rotation mismatch for {result.get('path', 'unknown')}: expected {expected_rotation}°, got {actual_rotation}°")
            return 0.0
    min_size = expected.get('min_size', 0)
    actual_size = result.get('file_size', 0)
    if min_size > 0 and actual_size < min_size:
        logger.warning(f"Video file too small for {result.get('path', 'unknown')}: expected >={min_size} bytes, got {actual_size} bytes")
        return 0.0
    logger.info(f"Video {result.get('path', 'unknown')}: rotation={actual_rotation}°, size={actual_size} bytes - PASS")
    return 1.0

def check_vlc_bgcone__ed0d0c08edef23629254f099c29e5e89(actual_config_path: str, expected: Dict) -> float:
    """
    Checks if VLC's background cone setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC config file
        expected: Dict containing 'bgcone_enabled' key (1 for enabled, 0 for disabled)

    Returns:
        1.0 if background cone setting matches expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_bgcone = expected['bgcone_enabled']
    if isinstance(expected_qt_bgcone, int):
        expected_qt_bgcone = str(expected_qt_bgcone)
    try:
        qt_bgcone = '1'
        for line in config_file.split('\n'):
            if 'qt-bgcone=' in line:
                qt_bgcone = line.split('=')[-1].strip()
        if qt_bgcone == expected_qt_bgcone:
            return 1.0
        else:
            logger.warning(f'Background cone mismatch - Expected: {expected_qt_bgcone}, Found: {qt_bgcone}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_screenshot_path__6c9e8f1d4a3b2e5f7890abcd12345678(actual_config_path: str, expected: dict, **options) -> float:
    """
    Checks if VLC's screenshot directory is set to the expected value.

    Args:
        actual_config_path: Path to the VLC config file
        expected: Expected rules dict containing 'screenshot_path'
        **options: Additional options

    Returns:
        1.0 if screenshot path matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_screenshot_path = expected['screenshot_path']
    try:
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-path' in line:
                current_path = line.split('=')[-1].strip()
                if current_path == expected_screenshot_path:
                    return 1.0
                else:
                    return 0.0
        return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_bitrate__e0916c40(result, expected, **options):
    """
    Check if video bitrate was correctly extracted.

    Args:
        result: Path to the bitrate.txt file
        expected: Dict with min and max bitrate in kbps
        **options: Additional options

    Returns:
        float: 1.0 if bitrate is valid, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        bitrate = int(content)
        min_bitrate = expected.get('min_bitrate_kbps', 0)
        max_bitrate = expected.get('max_bitrate_kbps', float('inf'))
        if min_bitrate <= bitrate <= max_bitrate:
            logger.info(f'Bitrate {bitrate} kbps is within valid range [{min_bitrate}, {max_bitrate}]')
            return 1.0
        else:
            logger.info(f'Bitrate {bitrate} kbps is outside valid range [{min_bitrate}, {max_bitrate}]')
            return 0.0
    except ValueError as e:
        logger.error(f'Could not parse bitrate as integer: {e}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking bitrate: {e}')
        return 0.0

def check_vlc_start_paused__a9847f6d(actual_config_path, rule):
    """
    Checks if VLC's start-paused setting is configured correctly.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_start_paused = rule['expected_start_paused']
    if isinstance(expected_start_paused, int):
        expected_start_paused = str(expected_start_paused)
    try:
        start_paused = '0'
        for line in config_file.split('\n'):
            if 'start-paused=' in line:
                start_paused = line.split('=')[-1].strip()
        if start_paused == expected_start_paused:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_paused__9bec8597(actual_status_path: str, expected) -> float:
    """
    Checks if VLC is currently paused.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Expected state ('paused')

    Returns:
        1.0 if VLC is paused, 0.0 otherwise
    """
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        status = tree.find('state').text
        logger.info(f'VLC Status: {status}')
        expected_state = expected if isinstance(expected, str) else expected.get('state', 'paused')
        if status == expected_state:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC pause status: {e}')
        return 0.0

def check_vlc_global_hotkey_stop__dfe711eae38b1244e43a6b27c52f3be0(actual_config_path, rule):
    """
    Checks if VLC's global hotkey for stop is set with a valid keyboard shortcut.

    This function not only checks if the global-key-stop is set, but also validates
    that the value is a valid keyboard shortcut (not just any arbitrary string).
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_global_key_stop = rule['expected_global_key_stop']
    if isinstance(expected_global_key_stop, int):
        expected_global_key_stop = str(expected_global_key_stop)
    try:
        global_key_stop = '0'
        for line in config_file.split('\n'):
            if 'global-key-stop=' in line:
                shortcut_value = line.split('=')[-1].strip()
                if shortcut_value != '' and is_valid_keyboard_shortcut(shortcut_value):
                    global_key_stop = '1'
                else:
                    global_key_stop = '0'
        if global_key_stop == expected_global_key_stop:
            return 1
        else:
            return 0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0

def check_qt_system_tray__629fd64a(actual_config_path, expected):
    """
    Check if VLC system tray icon is enabled.
    This checks the qt-system-tray setting in VLC config.

    Args:
        actual_config_path: Path to vlcrc config file
        expected: Expected value (0=disabled, 1=enabled)

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_system_tray = expected.get('expected_qt_system_tray')
    if isinstance(expected_qt_system_tray, int):
        expected_qt_system_tray = str(expected_qt_system_tray)
    try:
        qt_system_tray = '1'
        for line in config_file.split('\n'):
            if 'qt-system-tray=' in line:
                qt_system_tray = line.split('=')[-1].strip()
        if qt_system_tray == expected_qt_system_tray:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_screenshot_format__bb7a94b9(actual_config_path, expected):
    """
    Checks if VLC's screenshot format is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_format = expected.get('screenshot_format')
    try:
        screenshot_format = 'png'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-format=' in line:
                screenshot_format = line.split('=')[-1].strip()
                break
        if screenshot_format == expected_format:
            return 1.0
        else:
            logger.warning(f'Screenshot format mismatch - Expected: {expected_format}, Found: {screenshot_format}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_snapshot_prefix__bdf11a45(result, expected, **options):
    """Check if VLC snapshot prefix matches expected value.

    Args:
        result: Current snapshot prefix from getter
        expected: Expected prefix from rules
        **options: Additional options

    Returns:
        1.0 if matches, 0.0 otherwise
    """
    expected_prefix = expected.get('prefix', 'interstellar-')
    if result == expected_prefix:
        return 1.0
    else:
        return 0.0

def check_vlc_loop_status__9e0433d6(result, expected, **options):
    """
    Compare VLC loop status against expected value.

    Args:
        result: Current loop status from getter (bool or None)
        expected: Expected loop status (bool)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC loop status result is None')
        return 0.0
    if result == expected:
        return 1.0
    else:
        logger.info(f'Loop status mismatch - Expected: {expected}, Got: {result}')
        return 0.0

def check_snapshot_path__f6cd7aaf(actual_config_path, rule):
    """
    Checks if VLC's snapshot-path setting is set to the expected directory.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_snapshot_path = rule['expected_snapshot_path']
    try:
        snapshot_path = ''
        for line in config_file.split('\n'):
            if 'snapshot-path=' in line:
                snapshot_path = line.split('=')[-1].strip()
        if snapshot_path == expected_snapshot_path:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_duration__ff469bc6(result, expected, **options):
    """
    Check if the video duration is within expected range.

    Args:
        result: Path to the duration.txt file
        expected: Dict with "min_duration" and "max_duration" keys
        **options: Additional options

    Returns:
        float: 1.0 if within range, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        duration = int(content)
        min_dur = expected.get('min_duration', 0)
        max_dur = expected.get('max_duration', float('inf'))
        if min_dur <= duration <= max_dur:
            logger.info(f'Duration {duration} is within range [{min_dur}, {max_dur}]')
            return 1.0
        else:
            logger.info(f'Duration {duration} is outside range [{min_dur}, {max_dur}]')
            return 0.0
    except ValueError as e:
        logger.error(f'Could not parse duration as integer: {e}')
        return 0.0
    except Exception as e:
        logger.error(f'Error reading file: {e}')
        return 0.0

def check_vlc_always_on_top__b1199d7f2afb5cd10588e7e1e6cdc076(result, expected, **options):
    """
    Checks if VLC always-on-top setting matches the expected state.

    Args:
        result: The current always-on-top status from VLC config (boolean)
        expected: The expected state from rules

    Returns:
        float: 1.0 if always-on-top state matches, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC always-on-top status is None')
        return 0.0
    expected_on_top = expected.get('on_top_enabled', True)
    if result == expected_on_top:
        logger.info(f'VLC always-on-top state matches expected: {expected_on_top}')
        return 1.0
    else:
        logger.warning(f'VLC always-on-top state mismatch - Expected: {expected_on_top}, Got: {result}')
        return 0.0

def check_vlc_video_transformed__5e99ab87(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly transformed with 180° rotation.

    For 180° rotation verification:
    1. Verify file was re-encoded (file size differs)
    2. Verify corner regions swapped positions (top-left ↔ bottom-right, top-right ↔ bottom-left)
    3. This proves the video was actually rotated 180°, not just any transformation

    Args:
        result: Combined output from ls, stat, ffprobe, and corner hash commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was correctly rotated 180°, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    remaining_after_source = probe_parts[1] if len(probe_parts) > 1 else ''
    corner_parts = remaining_after_source.split('---SRC_TL---')
    source_video_info = corner_parts[0].strip() if len(corner_parts) > 0 else ''
    src_tl = src_tr = src_bl = src_br = ''
    out_tl = out_tr = out_bl = out_br = ''
    if len(corner_parts) > 1:
        remaining = corner_parts[1]
        parts_tr = remaining.split('---SRC_TR---')
        src_tl = parts_tr[0].strip() if len(parts_tr) > 0 else ''
        if len(parts_tr) > 1:
            parts_bl = parts_tr[1].split('---SRC_BL---')
            src_tr = parts_bl[0].strip() if len(parts_bl) > 0 else ''
            if len(parts_bl) > 1:
                parts_br = parts_bl[1].split('---SRC_BR---')
                src_bl = parts_br[0].strip() if len(parts_br) > 0 else ''
                if len(parts_br) > 1:
                    parts_out_tl = parts_br[1].split('---OUT_TL---')
                    src_br = parts_out_tl[0].strip() if len(parts_out_tl) > 0 else ''
                    if len(parts_out_tl) > 1:
                        parts_out_tr = parts_out_tl[1].split('---OUT_TR---')
                        out_tl = parts_out_tr[0].strip() if len(parts_out_tr) > 0 else ''
                        if len(parts_out_tr) > 1:
                            parts_out_bl = parts_out_tr[1].split('---OUT_BL---')
                            out_tr = parts_out_bl[0].strip() if len(parts_out_bl) > 0 else ''
                            if len(parts_out_bl) > 1:
                                parts_out_br = parts_out_bl[1].split('---OUT_BR---')
                                out_bl = parts_out_br[0].strip() if len(parts_out_br) > 0 else ''
                                if len(parts_out_br) > 1:
                                    out_br = parts_out_br[1].strip()
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0

    def extract_hash(framehash_str):
        if not framehash_str or 'ERROR' in framehash_str:
            return ''
        parts = framehash_str.split(',')
        return parts[-1].strip() if len(parts) >= 3 else framehash_str.strip()
    src_tl_hash = extract_hash(src_tl)
    src_tr_hash = extract_hash(src_tr)
    src_bl_hash = extract_hash(src_bl)
    src_br_hash = extract_hash(src_br)
    out_tl_hash = extract_hash(out_tl)
    out_tr_hash = extract_hash(out_tr)
    out_bl_hash = extract_hash(out_bl)
    out_br_hash = extract_hash(out_br)
    if output_filesize is None or source_filesize is None:
        logger.warning('Cannot verify file sizes')
        return 0.0
    size_diff = abs(output_filesize - source_filesize)
    size_diff_percent = size_diff / source_filesize * 100 if source_filesize > 0 else 0
    if size_diff < 1024:
        logger.warning(f'File size difference too small ({size_diff} bytes) - likely not re-encoded')
        return 0.0
    logger.info(f'File size changed: {source_filesize} -> {output_filesize} ({size_diff_percent:.2f}%)')
    if not all([src_tl_hash, src_tr_hash, src_bl_hash, src_br_hash, out_tl_hash, out_tr_hash, out_bl_hash, out_br_hash]):
        logger.warning('Missing corner hash data - cannot verify 180° rotation')
        logger.warning(f'Source corners: TL={bool(src_tl_hash)}, TR={bool(src_tr_hash)}, BL={bool(src_bl_hash)}, BR={bool(src_br_hash)}')
        logger.warning(f'Output corners: TL={bool(out_tl_hash)}, TR={bool(out_tr_hash)}, BL={bool(out_bl_hash)}, BR={bool(out_br_hash)}')
        return 0.0
    all_hashes = [src_tl_hash, src_tr_hash, src_bl_hash, src_br_hash, out_tl_hash, out_tr_hash, out_bl_hash, out_br_hash]
    if any((len(h) < 16 for h in all_hashes)):
        logger.warning('Invalid hash data - hashes too short')
        return 0.0
    tl_br_match = src_tl_hash == out_br_hash
    tr_bl_match = src_tr_hash == out_bl_hash
    bl_tr_match = src_bl_hash == out_tr_hash
    br_tl_match = src_br_hash == out_tl_hash
    logger.info(f'180° rotation corner verification:')
    logger.info(f'  Source TL → Output BR: {tl_br_match} (src={src_tl_hash[:16]}..., out={out_br_hash[:16]}...)')
    logger.info(f'  Source TR → Output BL: {tr_bl_match} (src={src_tr_hash[:16]}..., out={out_bl_hash[:16]}...)')
    logger.info(f'  Source BL → Output TR: {bl_tr_match} (src={src_bl_hash[:16]}..., out={out_tr_hash[:16]}...)')
    logger.info(f'  Source BR → Output TL: {br_tl_match} (src={src_br_hash[:16]}..., out={out_tl_hash[:16]}...)')
    matches = [tl_br_match, tr_bl_match, bl_tr_match, br_tl_match]
    matches_count = sum(matches)
    if matches_count == 4:
        logger.info(f'✓ Video correctly rotated 180°: all 4 corners swapped positions')
        return 1.0
    elif matches_count >= 3:
        logger.info(f'⚠ Video likely rotated 180°: {matches_count}/4 corners match (allowing compression tolerance)')
        return 1.0
    else:
        logger.warning(f'✗ Video NOT correctly rotated 180°: only {matches_count}/4 corners match')
        logger.warning(f'   This suggests a different transformation (not 180° rotation)')
        return 0.0

def check_qt_video_autoresize__6947249d(actual_config_path, expected):
    """
    Check if VLC is configured to auto-resize the video window.
    This checks the qt-video-autoresize setting in VLC config.

    Args:
        actual_config_path: Path to vlcrc config file
        expected: Expected value (0=disabled, 1=enabled)

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_video_autoresize = expected.get('expected_qt_video_autoresize')
    if isinstance(expected_qt_video_autoresize, int):
        expected_qt_video_autoresize = str(expected_qt_video_autoresize)
    try:
        qt_video_autoresize = '1'
        for line in config_file.split('\n'):
            if 'qt-video-autoresize=' in line:
                qt_video_autoresize = line.split('=')[-1].strip()
        if qt_video_autoresize == expected_qt_video_autoresize:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_muted__95775285be787fb106c6c57a61517d2a(actual_status_path: str, expected: Dict) -> float:
    """
    Checks if VLC audio is muted.

    This function handles VLC mute state detection by checking for available indicators:

    1. PRIORITY 1: Explicit mute flag (if VLC version exposes it)
       - Field name may be 'mute', 'muted', 'audio_mute', or 'is_muted'
       - If present, this is the definitive indicator

    2. PRIORITY 2: Volume level inference (fallback for all VLC versions)
       - VLC's mute button implementation sets reported volume to 0
       - This is reliable because:
         a) Clicking mute button always sets volume=0 in HTTP status
         b) Unmuting restores previous volume level
         c) While theoretically volume and mute are independent, VLC's UI
            implementation couples them in the HTTP interface
       - Note: If agent manually drags volume slider to 0 without clicking
         mute button, this will pass (acceptable behavior - audio is still silent)

    The dual-check approach (explicit flag first, volume fallback) ensures:
    - Correct verification across VLC versions
    - Handles both explicit mute and volume-based mute detection
    - Practical over theoretical purity (silence achieved = task success)

    Args:
        actual_status_path: Path to VLC status file (JSON format from getter)
        expected: Dict with expected_muted (true/false)

    Returns:
        1.0 if mute status matches expected, 0.0 otherwise
    """
    if not actual_status_path:
        logger.error('No VLC status path provided')
        return 0.0
    try:
        with open(actual_status_path, 'r') as file:
            status_data = json.load(file)
        actual_muted = False
        if 'muted_explicit' in status_data:
            actual_muted = bool(status_data['muted_explicit'])
            logger.info(f'Using explicit mute flag from VLC: {actual_muted}')
        else:
            volume = status_data.get('volume', 256)
            actual_muted = volume == 0
            logger.info(f'Inferred mute state from volume level: volume={volume}, muted={actual_muted}')
        expected_muted = expected.get('expected_muted', True)
        logger.info(f"VLC Mute Check - Actual: {actual_muted}, Expected: {expected_muted}, Result: {('PASS' if actual_muted == expected_muted else 'FAIL')}")
        return 1.0 if actual_muted == expected_muted else 0.0
    except FileNotFoundError:
        logger.error(f'VLC status file not found: {actual_status_path}')
        return 0.0
    except json.JSONDecodeError as e:
        logger.error(f'Failed to parse JSON status file: {e}')
        return _check_vlc_muted_xml_fallback(actual_status_path, expected)
    except Exception as e:
        logger.error(f'Error checking VLC muted status: {e}')
        return 0.0

def check_snapshot_format__a060a2ea(result, expected, **options):
    """Check if snapshot file format matches expected.

    Args:
        result: File format from getter
        expected: Expected format from rules
        **options: Additional options

    Returns:
        1.0 if format matches, 0.0 otherwise
    """
    expected_format = expected.get('format', 'PNG')
    if result == expected_format:
        return 1.0
    else:
        return 0.0

def check_vlc_always_on_top__245c3f85(result, expected, **options):
    """
    Check if VLC's always-on-top mode is enabled as expected.

    Args:
        result: Boolean indicating if always-on-top is enabled
        expected: Expected configuration
        **options: Additional options

    Returns:
        float: 1.0 if state matches expected, 0.0 otherwise
    """
    should_be_on_top = expected.get('should_be_on_top', True)
    if should_be_on_top:
        return 1.0 if result else 0.0
    else:
        return 0.0 if result else 1.0

def check_vlc_screenshot_saved__27395ab4949ce7b3fcdc03fed506b8b7(result, expected, **options):
    """
    Check if VLC screenshot was saved with correct filename and format.

    Args:
        result: Screenshot info from getter
        expected: Expected filename and properties
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        logger.info('Screenshot file does not exist')
        return 0.0
    expected_filename = expected.get('filename', '')
    actual_filename = result.get('filename', '')
    if actual_filename == expected_filename:
        score += 0.3
    else:
        logger.warning(f"Filename mismatch: expected '{expected_filename}', got '{actual_filename}'")
    if result.get('is_png', False):
        score += 0.3
    else:
        logger.warning('File is not a valid PNG image')
    logger.info(f'VLC screenshot check score: {score}')
    return score

def check_snapshot_count__8588412f(result, expected, **options):
    """Check if snapshot file count meets minimum requirement.

    Args:
        result: File count from getter
        expected: Expected count from rules
        **options: Additional options

    Returns:
        1.0 if count >= minimum, 0.0 otherwise
    """
    min_count = expected.get('min_count', 3)
    if result >= min_count:
        return 1.0
    else:
        return 0.0

def check_wav_conversion__580377a3f955e45d41d67d84cdd5fa88(result, expected, **options):
    """
    Verify WAV audio file was created correctly from video.
    """
    score = 0.0
    if result.get('exists', False):
        score += 0.35
        logger.info(f"WAV file exists: {result.get('filename', 'unknown')}")
    else:
        logger.warning('WAV file does not exist')
        return score
    expected_filename = expected.get('filename', '')
    if result.get('filename', '') == expected_filename:
        score += 0.25
        logger.info(f'Filename correct: {expected_filename}')
    else:
        logger.warning(f"Filename incorrect - Expected: {expected_filename}, Got: {result.get('filename', '')}")
    if result.get('is_wav', False):
        score += 0.25
        logger.info(f"Format is WAV: {result.get('mime_type', 'unknown')}")
    else:
        logger.warning(f"Not WAV format: {result.get('mime_type', 'unknown')}")
    min_size = expected.get('min_size', 100000)
    if result.get('size', 0) >= min_size:
        score += 0.15
        logger.info(f"File size adequate: {result.get('size', 0)} bytes")
    else:
        logger.warning(f"File size too small: {result.get('size', 0)} bytes (min: {min_size})")
    return score

def check_vlc_minimal_view__99af05cec0829d789ac2d7bf7abe0481(actual_config_path: str, expected: Dict) -> float:
    """
    Checks if VLC's minimal view mode is set to the expected value.

    Args:
        actual_config_path: Path to the VLC config file
        expected: Dict containing 'minimal_view_enabled' key (1 for enabled, 0 for disabled)

    Returns:
        1.0 if minimal view setting matches expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_minimal_view = expected['minimal_view_enabled']
    if isinstance(expected_minimal_view, int):
        expected_minimal_view = str(expected_minimal_view)
    try:
        qt_minimal_view = '0'
        for line in config_file.split('\n'):
            if 'qt-minimal-view=' in line:
                qt_minimal_view = line.split('=')[-1].strip()
        if qt_minimal_view == expected_minimal_view:
            return 1.0
        else:
            logger.warning(f'Minimal view mismatch - Expected: {expected_minimal_view}, Found: {qt_minimal_view}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_loop_mode__6d45f22f(actual_config_path, rule):
    """
    Checks if VLC's loop/repeat setting is configured correctly.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_loop = rule['expected_loop']
    if isinstance(expected_loop, int):
        expected_loop = str(expected_loop)
    try:
        loop = '0'
        for line in config_file.split('\n'):
            if line.startswith('loop=') or ' loop=' in line:
                loop = line.split('=')[-1].strip()
        if loop == expected_loop:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_playback_rate__a884bc65(result, expected, **options):
    """
    Compare VLC playback rate against expected value.

    Args:
        result: Current playback rate from getter (float or None)
        expected: Expected playback rate (float)
        **options: Additional options (tolerance for approximate match)

    Returns:
        float: 1.0 if match (within tolerance), 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC playback rate result is None')
        return 0.0
    tolerance = options.get('tolerance', 0.05)
    if abs(result - expected) <= tolerance:
        return 1.0
    else:
        logger.info(f'Playback rate mismatch - Expected: {expected}, Got: {result}, Tolerance: {tolerance}')
        return 0.0

def check_vlc_playing_file__24795e62(result, expected, **options):
    """Check if VLC is playing and verify the filename.

    Args:
        result: Path to VLC status XML file
        expected: Expected rules dict with 'file_name' key
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not os.path.exists(result):
        logger.warning('VLC status file not found')
        return 0.0
    with open(result, 'rb') as file:
        actual_status = file.read().decode('utf-8')
    try:
        tree = ElementTree.fromstring(actual_status)
        state = tree.find('state')
        if state is None or state.text != 'playing':
            logger.info(f"VLC is not playing (state: {(state.text if state else 'unknown')})")
            return 0.0
        expected_filename = expected.get('file_name')
        if expected_filename:
            file_paths = ['information/category[@name="meta"]/info[@name="filename"]', 'information/category[@name="meta"]/info[@name="title"]']
            file_info = None
            for path in file_paths:
                element = tree.find(path)
                if element is not None and element.text:
                    file_info = element.text
                    break
            if file_info:
                actual_basename = os.path.basename(file_info)
                if actual_basename == expected_filename or file_info.endswith(expected_filename):
                    return 1.0
                logger.warning(f'File name mismatch - Expected: {expected_filename}, Found: {file_info}')
                return 0.0
            else:
                logger.warning('Could not find file information in VLC status')
                return 0.0
        return 1.0
    except Exception as e:
        logger.error(f'Error parsing VLC status: {e}')
        return 0.0

def check_vlc_snapshot_directory__b9846f28(actual_config_path, expected):
    """
    Checks if VLC's snapshot directory setting is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected configuration (from rules dict)

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_snapshot_path = expected['expected_snapshot_path']
    try:
        snapshot_path = None
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'snapshot-path=' in line:
                snapshot_path = line.split('=')[-1].strip()
                break
        if snapshot_path == expected_snapshot_path:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_stopped__b038ba5d0dbf1c0974f6a77a9290fde5(actual_status_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC is currently stopped (not playing or paused).
    """
    with open(actual_status_path, 'rb') as file:
        actual_status = file.read().decode('utf-8')
    tree = ElementTree.fromstring(actual_status)
    status = tree.find('state').text
    logger.info(f'VLC Status: {status}')
    if status == 'stopped':
        return 1.0
    else:
        logger.warning(f'VLC is not stopped. Current state: {status}')
        return 0.0

def check_vlc_video_transformed__bdd54294(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly rotated by the specified angle.

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename', 'source_filename', and 'expected_rotation' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was rotated correctly, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    expected_rotation = expected.get('expected_rotation', None)
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    transformation_detected = False
    if output_filesize is not None and source_filesize is not None:
        if output_filesize != source_filesize:
            logger.info(f'File size changed: {source_filesize} -> {output_filesize}')
            transformation_detected = True
        else:
            logger.warning('File sizes identical - likely renamed/copied, not transformed')
    output_rotation = output_info.get('rotate', '0')
    source_rotation = source_info.get('rotate', '0') if source_info else '0'
    try:
        output_rot_value = int(output_rotation) if output_rotation else 0
        source_rot_value = int(source_rotation) if source_rotation else 0
    except (ValueError, TypeError):
        output_rot_value = 0
        source_rot_value = 0
    if expected_rotation:
        expected_rot_value = int(expected_rotation)
        rotation_matches = False
        if expected_rot_value == 270:
            calculated_rotation = (source_rot_value + expected_rot_value) % 360
            if output_rot_value == 270 or output_rot_value == -90:
                logger.info(f'Rotation metadata matches expected: {output_rot_value} degrees')
                rotation_matches = True
            elif output_rot_value == calculated_rotation:
                logger.info(f'Rotation calculated correctly: {source_rot_value} + {expected_rot_value} = {output_rot_value}')
                rotation_matches = True
        elif expected_rot_value == -90:
            if output_rot_value == -90 or output_rot_value == 270:
                logger.info(f'Rotation metadata matches expected: {output_rot_value} degrees')
                rotation_matches = True
        else:
            calculated_rotation = (source_rot_value + expected_rot_value) % 360
            if output_rot_value == expected_rot_value or output_rot_value == calculated_rotation:
                logger.info(f'Rotation metadata matches expected: {output_rot_value} degrees')
                rotation_matches = True
        if source_info and (expected_rot_value == 270 or expected_rot_value == -90):
            output_width = output_info.get('width', '')
            output_height = output_info.get('height', '')
            source_width = source_info.get('width', '')
            source_height = source_info.get('height', '')
            if output_width == source_height and output_height == source_width and (output_width != output_height):
                logger.info(f'Dimensions swapped for 270° rotation: {source_width}x{source_height} -> {output_width}x{output_height}')
                rotation_matches = True
        if rotation_matches:
            logger.info(f'Video rotation verified: {expected_rot_value} degrees')
            return 1.0
        else:
            logger.warning(f'Rotation does not match expected {expected_rot_value} degrees. Output rotation: {output_rot_value}')
            return 0.0
    else:
        if output_rotation != source_rotation:
            logger.info(f'Rotation changed: {source_rotation} -> {output_rotation}')
            transformation_detected = True
        if source_info:
            output_width = output_info.get('width', '')
            output_height = output_info.get('height', '')
            source_width = source_info.get('width', '')
            source_height = source_info.get('height', '')
            if output_width == source_height and output_height == source_width and (output_width != output_height):
                logger.info(f'Dimensions swapped: {source_width}x{source_height} -> {output_width}x{output_height}')
                transformation_detected = True
        if transformation_detected:
            logger.info('Video transformation verified')
            return 1.0
        else:
            logger.warning('No clear evidence of video transformation')
            return 0.0

def check_vlc_repeat_mode__e611293b(actual_config_path, expected, **options):
    """Check if VLC repeat/loop mode is configured correctly.

    Args:
        actual_config_path: Path to VLC config file
        expected: Dictionary with 'expected_repeat_mode' key (0=off, 1=repeat one, 2=repeat all)
        **options: Additional options

    Returns:
        float: 1.0 if repeat mode matches expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_repeat = expected.get('expected_repeat_mode', '0')
    if isinstance(expected_repeat, int):
        expected_repeat = str(expected_repeat)
    try:
        repeat_mode = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'qt-continue=' in line:
                repeat_mode = line.split('=')[-1].strip()
                break
            elif 'loop=' in line:
                value = line.split('=')[-1].strip()
                if value == '1':
                    repeat_mode = '2'
                break
        logger.info(f'VLC repeat mode: {repeat_mode}, expected: {expected_repeat}')
        if repeat_mode == expected_repeat:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_video_metadata_extracted__bc3a25a7baab15992708f46f1d51e584(result, expected, **options):
    """
    Check if video metadata was properly extracted.

    Args:
        result: Dict with metadata status from getter
        expected: Dict with expected values (should_have_metadata: bool)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    should_have_metadata = expected.get('should_have_metadata', True)
    if not should_have_metadata and (not result.get('exists', False)):
        return 1.0
    if not should_have_metadata and result.get('exists', False):
        return 0.0
    score = 0.0
    if result.get('exists', False):
        score += 0.4
    else:
        return 0.0
    if result.get('has_errors', False):
        return score * 0.5
    if result.get('has_duration', False) and result.get('duration_value') is not None:
        score += 0.3
    if result.get('has_resolution', False) and result.get('resolution_value') is not None:
        score += 0.3
    return score

def check_vlc_max_volume__dea0bbbbee03e4923af38de1331e256f(actual_config_path: str, expected: Dict) -> float:
    """
    Checks if VLC's maximum volume setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC config file
        expected: Dict containing 'max_volume' key (integer percentage, e.g., 200 for 200%)

    Returns:
        1.0 if maximum volume setting matches expected, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_max_volume = expected['max_volume']
    if isinstance(expected_qt_max_volume, int):
        expected_qt_max_volume = str(expected_qt_max_volume)
    try:
        qt_max_volume = '125'
        for line in config_file.split('\n'):
            if 'qt-max-volume=' in line:
                qt_max_volume = line.split('=')[-1].strip()
        if qt_max_volume == expected_qt_max_volume:
            return 1.0
        else:
            logger.warning(f'Maximum volume mismatch - Expected: {expected_qt_max_volume}, Found: {qt_max_volume}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_loop__869e58e4ccfbeafe3d06c94499c355f5(actual_status_path: str, expected: Dict[str, bool]) -> float:
    """
    Checks if VLC loop/repeat mode is enabled.

    Args:
        actual_status_path: Path to VLC status XML file
        expected: Dict with expected loop state (should contain expected_loop=true)

    Returns:
        1.0 if loop matches expected state, 0.0 otherwise
    """
    if not actual_status_path:
        logger.error('No VLC status path provided')
        return 0.0
    try:
        with open(actual_status_path, 'rb') as file:
            actual_status = file.read().decode('utf-8')
        tree = ElementTree.fromstring(actual_status)
        loop_element = tree.find('loop')
        if loop_element is not None:
            actual_loop = loop_element.text.lower() == 'true'
        else:
            actual_loop = False
        expected_loop = expected.get('expected_loop', True)
        logger.info(f'VLC Loop Status: {actual_loop}, Expected: {expected_loop}')
        if actual_loop == expected_loop:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'Error checking VLC loop state: {e}')
        return 0.0

def check_vlc_loop_mode__03941735(actual_config_path, expected):
    """
    Checks if VLC's loop/repeat mode is enabled.

    Args:
        actual_config_path: Path to VLC config file
        expected: Expected loop mode setting (e.g., {"loop_enabled": "1"})

    Returns:
        1.0 if loop mode matches expected, 0.0 otherwise
    """
    try:
        with open(actual_config_path, 'rb') as file:
            config_file = file.read().decode('utf-8')
        expected_loop = expected.get('loop_enabled', '1')
        if isinstance(expected_loop, int):
            expected_loop = str(expected_loop)
        loop_mode = '0'
        for line in config_file.split('\n'):
            if 'loop=' in line and (not line.startswith('#')):
                loop_mode = line.split('=')[-1].strip()
                break
        if loop_mode == expected_loop:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_wallpaper_is_snapshot__b36f4827(result, expected, **options):
    """Check if wallpaper is a snapshot file.

    Args:
        result: Boolean from getter
        expected: Expected value from rules
        **options: Additional options

    Returns:
        Score based on partial completion
    """
    expected_is_snapshot = expected.get('is_snapshot', True)
    if result == expected_is_snapshot:
        return 1.0
    else:
        return 0.0

def check_srt_file_and_video__ed96ceb6(result, expected, **options):
    """
    Verify that subtitle file was extracted and video has no embedded subtitles.

    Args:
        result: Dict with 'srt_path' and 'has_subtitles' keys
        expected: Expected rules
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    score = 0.0
    srt_path = result.get('srt_path')
    has_subtitles_in_video = result.get('has_subtitles', True)
    if srt_path and os.path.exists(srt_path):
        if os.path.getsize(srt_path) > 0:
            try:
                subs = pysrt.open(srt_path)
                if len(subs) > 0:
                    score += 0.6
            except:
                pass
    if not has_subtitles_in_video:
        score += 0.4
    return score

def check_vlc_muted__eef15ca1487d76cdd0b6405615ef653d(result, expected, **options):
    """
    Checks if VLC audio is muted by verifying volume is 0.

    IMPORTANT LIMITATION DOCUMENTED:
    VLC's HTTP interface (status.xml) does NOT expose a separate mute boolean flag.
    The HTTP API only provides volume level (0-256+), not the internal mute state flag.

    When the user clicks VLC's mute button (or presses 'M' key), VLC internally:
    1. Sets volume to 0
    2. Stores the previous volume level for unmute
    However, the HTTP interface only sees the volume=0, not the mute flag itself.

    VERIFICATION BEHAVIOR:
    This metric checks if volume == 0 to determine mute state. This means:
    - User clicks mute button → volume becomes 0 → PASS ✓
    - User manually sets volume to 0 → volume becomes 0 → PASS ✓

    Both actions are functionally equivalent from the HTTP API's perspective and
    both achieve the task goal of "muting audio in VLC". The instruction asks to
    "mute the audio", which can be accomplished by either method, and both result
    in no audio output (volume=0).

    RATIONALE FOR volume == 0 (not volume < 5):
    - Volume == 0 is the precise indicator that audio is muted/silenced
    - Using threshold < 5 would incorrectly accept volume=1,2,3,4 as "muted"
    - VLC volume is an integer, so 0 is exact and unambiguous
    - If VLC volume is 1-4, audio is still audible (not muted)

    Args:
        result: Dict containing audio state with 'volume' and 'vlc_responsive' keys
        expected: The expected mute state from rules (dict with 'is_muted' key)

    Returns:
        float: 1.0 if mute state matches expected, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC audio state is None - VLC may not be running or responsive')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'Expected dict from getter, got {type(result)}')
        return 0.0
    vlc_responsive = result.get('vlc_responsive', False)
    if not vlc_responsive:
        logger.error('VLC is not responsive - cannot verify mute state')
        return 0.0
    volume = result.get('volume', None)
    if volume is None:
        logger.warning('Volume information not available from VLC')
        return 0.0
    expected_muted = expected.get('is_muted', True)
    is_muted = volume == 0
    if is_muted == expected_muted:
        logger.info(f'VLC mute state matches expected: {expected_muted} (volume: {volume}, is_muted: {is_muted})')
        return 1.0
    else:
        logger.warning(f'VLC mute state mismatch - Expected muted: {expected_muted}, Got volume: {volume}, computed is_muted: {is_muted}')
        return 0.0

def check_video_framerate__8c156dd1(result, expected, **options):
    """
    Check if video framerate was correctly extracted.

    Args:
        result: Path to the framerate.txt file
        expected: Dict with min_fps and max_fps
        **options: Additional options

    Returns:
        float: 1.0 if framerate is valid, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        fps = float(content)
        min_fps = expected.get('min_fps', 0.0)
        max_fps = expected.get('max_fps', float('inf'))
        if min_fps <= fps <= max_fps:
            logger.info(f'Framerate {fps} is within valid range [{min_fps}, {max_fps}]')
            return 1.0
        else:
            logger.info(f'Framerate {fps} is outside valid range [{min_fps}, {max_fps}]')
            return 0.0
    except ValueError as e:
        logger.error(f'Could not parse framerate as number: {e}')
        return 0.0
    except Exception as e:
        logger.error(f'Error checking framerate: {e}')
        return 0.0

def check_vlc_audio_muted__b5c6209b(result, expected, **options):
    """
    Compare VLC audio mute status against expected value.

    Args:
        result: Current mute status from getter (bool or None)
        expected: Expected mute status (bool)
        **options: Additional options

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC audio mute status result is None')
        return 0.0
    if result == expected:
        return 1.0
    else:
        logger.info(f'Audio mute status mismatch - Expected: {expected}, Got: {result}')
        return 0.0

def check_vlc_qt_video_autoresize__06c10f73(actual_config_path, expected):
    """Check VLC video autoresize setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_autoresize = expected.get('expected_video_autoresize')
    if isinstance(expected_autoresize, int):
        expected_autoresize = str(expected_autoresize)
    try:
        video_autoresize = '1'
        for line in config_file.split('\n'):
            if 'qt-video-autoresize=' in line and (not line.strip().startswith('#')):
                video_autoresize = line.split('=')[-1].strip()
        if video_autoresize == expected_autoresize:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_start_paused__ae1f301e(actual_config_path, expected):
    """
    Checks if VLC's start-paused setting is set to the expected value.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_paused = expected.get('start_paused')
    if isinstance(expected_paused, int):
        expected_paused = str(expected_paused)
    try:
        start_paused = '0'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'start-paused=' in line:
                start_paused = line.split('=')[-1].strip()
                break
        if start_paused == expected_paused:
            return 1.0
        else:
            logger.warning(f'Start-paused mismatch - Expected: {expected_paused}, Found: {start_paused}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_loop_setting__0435b04a(actual_config_path, expected):
    """Check VLC loop/repeat setting."""
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_loop = expected.get('expected_loop')
    if isinstance(expected_loop, int):
        expected_loop = str(expected_loop)
    try:
        loop_value = '0'
        for line in config_file.split('\n'):
            if 'loop=' in line and (not line.strip().startswith('#')):
                loop_value = line.split('=')[-1].strip()
        if loop_value == expected_loop:
            return 1.0
        else:
            return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_audio_duration__ee9c4304f47d297d44dedaad1e2983d6(result, expected, **options):
    """
    Check if audio file exists and has reasonable duration.

    Args:
        result: dict from getter with 'exists', 'duration', 'size' keys
        expected: dict with 'min_duration' and optionally 'max_duration'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not result.get('exists', False):
        logger.info('File does not exist')
        return 0.0
    score = 0.0
    score += 0.5
    min_duration = expected.get('min_duration', 0)
    max_duration = expected.get('max_duration', float('inf'))
    actual_duration = result.get('duration', 0)
    if min_duration <= actual_duration <= max_duration:
        score += 0.5
        logger.info(f'Duration {actual_duration}s is within range [{min_duration}, {max_duration}]')
    else:
        logger.warning(f'Duration {actual_duration}s is outside range [{min_duration}, {max_duration}]')
    return score

def check_vlc_loop_mode__bd908fa8(result, expected, **options):
    """
    Check if VLC loop/repeat mode is enabled as expected.

    Args:
        result: Dictionary with loop_enabled and repeat_enabled flags
        expected: Expected configuration
        **options: Additional options

    Returns:
        float: 1.0 if loop or repeat is enabled, 0.0 otherwise
    """
    if not result or not isinstance(result, dict):
        return 0.0
    loop_should_be_enabled = expected.get('loop_should_be_enabled', True)
    is_loop_enabled = result.get('loop_enabled', False) or result.get('repeat_enabled', False)
    if loop_should_be_enabled:
        return 1.0 if is_loop_enabled else 0.0
    else:
        return 0.0 if is_loop_enabled else 1.0

def check_audio_file_created__5d993657(result, expected, **options):
    """Check if audio file was created with minimum size.

    Args:
        result: dict with 'exists', 'size', 'path' keys
        expected: dict with 'min_size' key

    Returns:
        float: 1.0 if file exists and meets size requirement, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.warning(f'Result is not a dict: {type(result)}')
        return 0.0
    exists = result.get('exists', False)
    size = result.get('size', 0)
    min_size = expected.get('min_size', 1000)
    if not exists:
        logger.info(f"Audio file does not exist at {result.get('path', 'unknown')}")
        return 0.0
    if size < min_size:
        logger.info(f'Audio file too small: {size} bytes < {min_size} bytes')
        return 0.0
    logger.info(f'Audio file created successfully: {size} bytes')
    return 1.0

def check_vlc_aspect_ratio__8a47ea01(result, expected, **options):
    """
    Check if VLC aspect ratio matches the expected value.

    Args:
        result: Current aspect ratio setting from getter
        expected: Expected configuration with 'aspect_ratio' key
        **options: Additional options

    Returns:
        float: 1.0 if aspect ratio matches expected, 0.0 otherwise
    """
    expected_ratio = expected.get('aspect_ratio')
    if expected_ratio is None:
        is_aspect_changed = result != 'default' and result != ''
        return 1.0 if is_aspect_changed else 0.0
    return 1.0 if result == expected_ratio else 0.0

def check_vlc_config_and_mp3_file__8f080098(result_state: Dict[str, Any], expected_state: Dict[str, Any], **options) -> float:
    """
    Check both VLC recording configuration and MP3 file creation.

    This metric verifies:
    1. VLC config has correct recording path (30% weight)
    2. MP3 file exists and has valid content (70% weight)

    Args:
        result_state: Dict with 'vlc_config' and 'file_info' keys
        expected_state: Dict with 'recording_path' and 'file_path' keys
        **options: Additional options

    Returns:
        float: Partial credit score between 0.0 and 1.0
    """
    logger.info(f'[DEBUG] check_vlc_config_and_mp3_file called')
    logger.info(f'[DEBUG] result_state: {result_state}')
    logger.info(f'[DEBUG] expected_state: {expected_state}')
    score = 0.0
    vlc_config = result_state.get('vlc_config', '')
    expected_recording_path = expected_state.get('recording_path', '')
    vlc_config_correct = False
    if vlc_config and expected_recording_path:
        try:
            for line in vlc_config.split('\n'):
                if line.startswith('#') or not line.strip():
                    continue
                if 'input-record-path' in line:
                    current_path = line.split('=')[-1].strip()
                    if current_path == expected_recording_path:
                        vlc_config_correct = True
                        logger.info(f'[DEBUG] VLC config correct: {current_path}')
                        break
                    else:
                        logger.info(f'[DEBUG] VLC config incorrect: {current_path} vs {expected_recording_path}')
                        break
        except Exception as e:
            logger.error(f'[DEBUG] Error checking VLC config: {e}')
    if vlc_config_correct:
        score += 0.3
        logger.info('[DEBUG] VLC config check passed: +0.3')
    else:
        logger.info('[DEBUG] VLC config check failed: +0.0')
    file_info = result_state.get('file_info', {})
    file_exists = file_info.get('exists', False)
    file_size = file_info.get('size', 0)
    if file_exists:
        score += 0.5
        logger.info('[DEBUG] MP3 file exists: +0.5')
        if file_size and file_size > 0:
            score += 0.2
            logger.info(f'[DEBUG] MP3 file has valid size ({file_size} bytes): +0.2')
        else:
            logger.info(f'[DEBUG] MP3 file is empty or size unknown: +0.0')
    else:
        logger.info('[DEBUG] MP3 file does not exist: +0.0')
    logger.info(f'[DEBUG] Final score: {score}')
    return score

def check_vlc_metadata_network_access__e85db42df496eec95b9a79d8fc8ac5e4(actual_config_path: str, rule: Dict[str, str]) -> float:
    """
    Checks if VLC's metadata network access setting is set to the expected value.
    The metadata-network-access setting controls whether VLC fetches metadata from the internet.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_metadata_access = rule['expected_metadata_network_access']
    if isinstance(expected_metadata_access, int):
        expected_metadata_access = str(expected_metadata_access)
    try:
        metadata_network_access = '1'
        for line in config_file.split('\n'):
            if line.startswith('#') or not line.strip():
                continue
            if 'metadata-network-access=' in line:
                metadata_network_access = line.split('=')[-1].strip()
                break
        if metadata_network_access == expected_metadata_access:
            return 1.0
        else:
            logger.warning(f'Metadata network access mismatch - Expected: {expected_metadata_access}, Found: {metadata_network_access}')
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_subtitle_autoload__aec93e6e(actual_config_path, rule):
    """
    Checks if VLC's subtitle autoload setting is configured correctly.
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_subtitle_autoload = rule['expected_subtitle_autoload']
    if isinstance(expected_subtitle_autoload, int):
        expected_subtitle_autoload = str(expected_subtitle_autoload)
    try:
        subtitle_autoload = '1'
        for line in config_file.split('\n'):
            if 'sub-autodetect-file=' in line:
                subtitle_autoload = line.split('=')[-1].strip()
        if subtitle_autoload == expected_subtitle_autoload:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_qt_privacy_ask__f3480833(actual_config_path, expected, **options):
    """
    Checks if VLC's privacy ask setting is set to the expected value.

    Args:
        actual_config_path: Path to the VLC configuration file
        expected: Expected value dict with "expected_qt_privacy_ask" key
        **options: Additional options

    Returns:
        float: 1.0 if setting matches, 0.0 otherwise
    """
    with open(actual_config_path, 'rb') as file:
        config_file = file.read().decode('utf-8')
    expected_qt_privacy_ask = expected['expected_qt_privacy_ask']
    if isinstance(expected_qt_privacy_ask, int):
        expected_qt_privacy_ask = str(expected_qt_privacy_ask)
    try:
        qt_privacy_ask = '1'
        for line in config_file.split('\n'):
            if 'qt-privacy-ask=' in line:
                qt_privacy_ask = line.split('=')[-1].strip()
        if qt_privacy_ask == expected_qt_privacy_ask:
            return 1.0
        else:
            return 0.0
    except FileNotFoundError:
        logger.error('VLC configuration file not found.')
        return 0.0
    except Exception as e:
        logger.error(f'An error occurred: {e}')
        return 0.0

def check_vlc_volume__c48d7656(result, expected, **options):
    """
    Compare VLC volume level against expected value.

    Args:
        result: Current volume level from getter (int or None)
        expected: Expected volume level (int)
        **options: Additional options (tolerance for approximate match)

    Returns:
        float: 1.0 if match (within tolerance), 0.0 otherwise
    """
    if result is None:
        logger.warning('VLC volume result is None')
        return 0.0
    tolerance = options.get('tolerance', 10)
    if abs(result - expected) <= tolerance:
        return 1.0
    else:
        logger.info(f'Volume mismatch - Expected: {expected}, Got: {result}, Tolerance: {tolerance}')
        return 0.0

def check_video_codec__05e4689a(result, expected, **options):
    """
    Check if the video codec matches expected values.

    Args:
        result: Path to the codec.txt file
        expected: Dict with "valid_codecs" list
        **options: Additional options

    Returns:
        float: 1.0 if codec is valid, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip().lower()
        valid_codecs = [c.lower() for c in expected.get('valid_codecs', [])]
        if content in valid_codecs:
            logger.info(f'Codec {content} is valid')
            return 1.0
        else:
            logger.info(f'Codec {content} not in valid list: {valid_codecs}')
            return 0.0
    except Exception as e:
        logger.error(f'Error reading file: {e}')
        return 0.0

def check_video_resolution__89b2d435(result, expected, **options):
    """
    Check if video resolution was correctly identified.

    Args:
        result: Path to the resolution.txt file
        expected: Dict with pattern and minimum dimensions
        **options: Additional options

    Returns:
        float: 1.0 if resolution is valid, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        with open(result, 'r') as f:
            content = f.read().strip()
        pattern = expected.get('pattern', '^\\d+x\\d+$')
        if not re.match(pattern, content):
            logger.info(f'Resolution format invalid: {content}')
            return 0.0
        parts = content.split('x')
        if len(parts) != 2:
            logger.info(f'Could not parse resolution: {content}')
            return 0.0
        width = int(parts[0])
        height = int(parts[1])
        min_width = expected.get('min_width', 0)
        min_height = expected.get('min_height', 0)
        if width < min_width or height < min_height:
            logger.info(f'Resolution {width}x{height} below minimum {min_width}x{min_height}')
            return 0.0
        logger.info(f'Resolution validated: {content}')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking resolution: {e}')
        return 0.0

def check_vlc_video_transformed__6e7fd09f(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if the video file exists and was properly transformed (rotated/flipped).

    Args:
        result: Combined output from ls, stat, and ffprobe commands
        expected: Dict with 'filename' and 'source_filename' keys
        **options: Additional options

    Returns:
        float: 1.0 if file exists and was transformed, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result is None')
        return 0.0
    expected_filename = expected.get('filename', '')
    parts = result.split('---SEPARATOR---')
    ls_output = parts[0] if len(parts) > 0 else result
    if 'No such file or directory' in ls_output or 'cannot access' in ls_output:
        logger.warning(f'Output file does not exist: {ls_output}')
        return 0.0
    if expected_filename and expected_filename not in ls_output:
        logger.warning(f"Output filename '{expected_filename}' not found in ls output")
        return 0.0
    logger.info(f'File exists: {expected_filename}')
    if len(parts) < 2:
        logger.warning('Missing stat/ffprobe output')
        return 0.5
    remaining = parts[1]
    filesize_parts = remaining.split('---FILESIZE_SEP---')
    if len(filesize_parts) < 2:
        logger.warning('Missing source file size')
        return 0.5
    output_filesize_str = filesize_parts[0].strip()
    remaining_after_filesize = filesize_parts[1]
    ffprobe_parts = remaining_after_filesize.split('---FFPROBE_OUT---')
    source_filesize_str = ffprobe_parts[0].strip() if len(ffprobe_parts) > 0 else ''
    if len(ffprobe_parts) < 2:
        logger.warning('Missing ffprobe output')
        return 0.5
    ffprobe_data = ffprobe_parts[1]
    probe_parts = ffprobe_data.split('---SOURCE_INFO---')
    output_video_info = probe_parts[0].strip() if len(probe_parts) > 0 else ''
    source_video_info = probe_parts[1].strip() if len(probe_parts) > 1 else ''
    try:
        output_filesize = int(output_filesize_str) if output_filesize_str.isdigit() else None
        source_filesize = int(source_filesize_str) if source_filesize_str.isdigit() else None
    except (ValueError, AttributeError):
        logger.warning(f'Could not parse file sizes')
        output_filesize = None
        source_filesize = None
    if not output_video_info or 'ERROR' in output_video_info:
        logger.warning('Could not read output video metadata')
        return 0.0

    def parse_video_info(info_str):
        info = {}
        lines = info_str.split('\n')
        for line in lines:
            line = line.strip()
            if '=' in line:
                (key, value) = line.split('=', 1)
                info[key] = value
        return info
    output_info = parse_video_info(output_video_info)
    source_info = parse_video_info(source_video_info) if source_video_info and 'ERROR' not in source_video_info else {}
    if 'width' not in output_info or 'height' not in output_info:
        logger.warning('Output file is not a valid video')
        return 0.0
    transformation_detected = False
    if output_filesize is not None and source_filesize is not None:
        if output_filesize != source_filesize:
            logger.info(f'File size changed: {source_filesize} -> {output_filesize}')
            transformation_detected = True
        else:
            logger.warning('File sizes identical - likely renamed/copied, not transformed')
    output_rotation = output_info.get('rotate', '0')
    source_rotation = source_info.get('rotate', '0') if source_info else '0'
    if output_rotation != source_rotation:
        logger.info(f'Rotation changed: {source_rotation} -> {output_rotation}')
        transformation_detected = True
    if source_info:
        output_width = output_info.get('width', '')
        output_height = output_info.get('height', '')
        source_width = source_info.get('width', '')
        source_height = source_info.get('height', '')
        if output_width == source_height and output_height == source_width and (output_width != output_height):
            logger.info(f'Dimensions swapped: {source_width}x{source_height} -> {output_width}x{output_height}')
            transformation_detected = True
    if transformation_detected:
        logger.info('Video transformation verified')
        return 1.0
    else:
        logger.warning('No clear evidence of video transformation')
        return 0.0

def check_video_snapshot__1391f174(result, expected, **options):
    """
    Check if video snapshot was properly created.

    Args:
        result: Path to the snapshot.jpg file
        expected: Dict with size and format requirements
        **options: Additional options

    Returns:
        float: 1.0 if snapshot created correctly, 0.0 otherwise
    """
    if result is None:
        logger.warning('Result file is None')
        return 0.0
    try:
        if not os.path.exists(result):
            logger.warning(f'File does not exist: {result}')
            return 0.0
        file_size = os.path.getsize(result)
        min_size = expected.get('min_size_bytes', 0)
        max_size = expected.get('max_size_bytes', float('inf'))
        if file_size < min_size or file_size > max_size:
            logger.info(f'File size {file_size} is outside range [{min_size}, {max_size}]')
            return 0.0
        if expected.get('check_jpeg', False):
            with open(result, 'rb') as f:
                header = f.read(3)
                if not header[0:2] == b'\xff\xd8':
                    logger.warning(f'File does not appear to be valid JPEG format')
                    return 0.0
        logger.info(f'Snapshot validated: {file_size} bytes')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking snapshot: {e}')
        return 0.0
