"""VeriGen generated judge functions.

Source: metrics.py
This module is auto-split from the original merged generated_tasks judge file.
"""

from PIL import Image
from PIL import Image, ImageChops
from PIL import Image, ImageFilter
from PIL import Image, ImageOps
from PIL import Image, ImageStat
from collections import Counter
from datetime import datetime
from datetime import datetime, timedelta
from datetime import time
from datetime import time, datetime
from desktop_env.evaluators.metrics.gimp import check_saturation_increase_and_structure_sim
from desktop_env.evaluators.metrics.gimp import structure_check_by_ssim
from desktop_env.evaluators.metrics.slides import check_strikethrough
from desktop_env.evaluators.metrics.slides import compare_pptx_files as _original_compare_pptx_files
from desktop_env.evaluators.metrics.utils import _match_record
from desktop_env.evaluators.metrics.utils import _match_value_to_rule as _match_pref
from desktop_env.evaluators.metrics.utils import compare_urls
from difflib import SequenceMatcher
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.text import WD_COLOR_INDEX
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.oxml import parse_xml
from docx.shared import Inches
from docx.shared import Inches, Pt
from docx.shared import Pt
from docx.shared import RGBColor
from docx.shared import RGBColor, Pt
from email import policy
from email.utils import parsedate_to_datetime
from io import BytesIO
from itertools import product
from lxml import etree
from lxml.cssselect import CSSSelector
from lxml.etree import _Element
from math import sqrt
from odf.draw import Frame
from odf.opendocument import load
from odf.style import Style, PageLayoutProperties
from openpyxl import Workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles import Color
from openpyxl.styles import PatternFill
from openpyxl.utils import column_index_from_string
from openpyxl.utils import get_column_letter
from openpyxl.utils import range_boundaries
from openpyxl.worksheet.worksheet import Worksheet
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.util import Inches
from pptx.util import Pt
from pypdf import PdfReader
from rapidfuzz import fuzz
from scipy import ndimage
from skimage.color import deltaE_ciede2000, rgb2lab
from skimage.metrics import structural_similarity as ssim
from typing import Any
from typing import Any, Dict
from typing import Any, Dict, List
from typing import Any, Dict, List, Tuple
from typing import Any, Dict, Optional
from typing import Any, Dict, Union
from typing import Any, List
from typing import Any, List, Dict
from typing import Any, List, Set
from typing import Any, Optional
from typing import Any, Optional, Dict
from typing import Dict
from typing import Dict, Any
from typing import Dict, Any, List
from typing import Dict, Any, List, Optional
from typing import Dict, Any, List, Tuple
from typing import Dict, Any, List, Union
from typing import Dict, Any, Optional
from typing import Dict, Any, Optional, List
from typing import Dict, Any, Optional, List, Tuple
from typing import Dict, Any, Optional, Tuple
from typing import Dict, Any, Pattern, Match
from typing import Dict, Any, Set
from typing import Dict, Any, Tuple
from typing import Dict, Any, Tuple, Optional
from typing import Dict, Any, Union
from typing import Dict, List
from typing import Dict, List, Any
from typing import Dict, List, Optional
from typing import Dict, List, Pattern
from typing import Dict, List, Union
from typing import Dict, Optional
from typing import Dict, Optional, Any
from typing import Dict, Tuple
from typing import Dict, Union
from typing import List
from typing import List, Any
from typing import List, Any, Dict
from typing import List, Any, Optional
from typing import List, Dict
from typing import List, Dict, Any
from typing import List, Dict, Any, Optional
from typing import List, Dict, Any, Optional, Tuple
from typing import List, Dict, Any, Tuple
from typing import List, Dict, Any, Union
from typing import List, Dict, Union, Pattern
from typing import List, Dict, Union, Pattern, Any
from typing import List, Optional
from typing import List, Optional, Any, Dict
from typing import List, Optional, Union
from typing import List, Pattern, Dict, Match
from typing import List, Tuple
from typing import List, Tuple, Dict
from typing import List, Union
from typing import Optional
from typing import Optional, Any
from typing import Optional, Dict
from typing import Optional, Dict, Any
from typing import Optional, Dict, Any, List
from typing import Optional, Dict, List
from typing import Optional, Tuple
from typing import Tuple, List, Dict, Any
from typing import Tuple, Optional
from typing import Union
from typing import Union, Any, TypeVar, Callable
from urllib.parse import urlparse
from urllib.parse import urlparse, parse_qs
from xml.etree import ElementTree
import PyPDF2
import ast
import csv
import cv2
import datetime
import difflib
import email
import fitz
import hashlib
import imagehash
import importlib.util
import io
import json
import logging
import lxml.etree
import numpy as np
import openpyxl
import openpyxl.utils
import operator
import os
import pytz
import re
import requests
import sqlite3
import subprocess
import sys
import tarfile
import tempfile
import time
import uuid
import xml.etree.ElementTree as ET
import zipfile

logger = logging.getLogger(__name__)

__all__ = ['is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_5', 'check_webext_manifest__bba937aa', 'check_chrome_experiments_not_contains__f72b4886483f6a090bcc6a28ba49acde', 'check_gdrive_files__da0d777d0a1e3c8d735dbb81a2e45a5c', 'check_gdrive_pdf_in_folder__b2131ee6', 'check_gdrive_pdf_in_folder__94eb0e23', 'check_exact_bookmarks__b53050a013ee52945e88dd9cbd9764a5', 'check_table_count_and_content__638d8a63', 'check_webext_manifest__dea4a5ea3d588b414bf151fee72b35b0', 'check_extension_manifest__1edda3bd4444a8fb1554a227ff178017', 'check_gdrive_pdf__3a6a494e', 'check_gdrive_pdf__9d82c3b2', 'is_expected_url_pattern_match__ac7e6c6cfb1c60f162b6d2bc92009d82', 'check_gdrive_files__45753efe', 'check_chrome_extension_manifest__6e0f538e8ca610f13ae0673161924889', 'is_expected_url_pattern_match_or__b070486d', 'check_chrome_ext_middle__484fc06534818b90d35ad638d0805c7c', 'check_extension_count__7a8c3f12', 'check_table_count_only__d0985f12', 'check_extension_description__221bd7ef80e48e21d62e7d38635cf26c', 'check_gdrive_pdf_in_folder__41049d6b', 'is_expected_search_query__e346411e', 'check_gdrive_pdf_in_folder__db0d2d11', 'check_csv_table_civil__2c0c636e', 'check_cookies_deleted__4730218c', 'is_extension_installed__72cc5a03', 'check_extension_installed__c548e99163664e4ea2d87d3cc236c650', 'check_chrome_setting__eec6beed', 'check_html_conversion_and_tab__6920e35a', 'check_gdrive_pdf__9e5c8fdc', 'check_chrome_setting__d84aae11', 'check_startup_new_tab__bd145bbb387b98da7ba508028ec58276', 'check_recreation_html_element__2aa2c046ed13ceae9537c9a8e566d558', 'check_vim_hlsearch__47d76eea6c988a4beb0cae6b4109cc24', 'check_recreation_html_element__6adae5a4637b133a46055994fbaa8dd4', 'check_webext_manifest__88e21052', 'is_expected_search_query__cbf92d65', 'check_backup_extension_files__8ece34a1', 'check_chrome_do_not_track__a858c66a', 'check_go_extension__303c34ab', 'check_table_top_right__199b082c', 'check_extension_exact_match__6410307fcb9f7a3e584e1435af0837eb', 'check_extension_description__2224641eff529090c9cc8ad62127e176', 'check_urls_present__92bc11a1', 'is_expected_search_query__e154bb96', 'check_chrome_experiments_contains__0a40109ea287ab7bbd8cd9175a7ce6a5', 'check_bookmarks_created__5641e4c0', 'check_pdf_and_tabs__a852a89d', 'check_csv_table_math__bdcc0312', 'is_extension_installed__7fabdee6', 'check_extension_count__0fb1bf36ddef8ab5a554de56ff7f0c3d', 'check_extension_dir_exists__42aa293803886f1362c8d3d5e33a2703', 'is_expected_url_pattern_match__9aac585d873d78f321c8f29733249832', 'check_chrome_extension_manifest__61167ddb747c3a88a31446374f5a958f', 'check_single_tab_url__c557d7de5c301473caee703726a81950', 'verify_tabs_with_state_change__58565672', 'check_table_count_only__e164fbd9', 'check_extension_name__de268faa9ad661d5adcd42e1b7f775e2', 'check_extension_folder_exists__ae6416e4', 'check_pdf_file_and_tab__a606126a', 'check_chrome_setting__0717aaed', 'check_gdrive_files__87199839', 'check_html_export_and_chrome__d6487d33', 'check_gdrive_pdf__7e0bdf48', 'is_expected_bookmarks__7a5a7856_aug_6_task_verify_0', 'check_gdrive_file__3ce93f6ce10147dcc2489b34874a9f3b', 'check_gdrive_eml_count__26d2516b888edf7c1b328cca7acaf9b7', 'check_webext_manifest__951cee96', 'check_gdrive_files__b0d15a73736689e26424d81f759f1528', 'is_expected_bookmarks__7a5a7856_f1b6_42a4_ade9_1ca81ca0f263_task_verify_0', 'check_html_file_and_tab__3cb1a587', 'check_website_filled__e7fe8e90', 'check_chrome_extension_manifest__70453fbf044a41274b6b7dc7e909fb11', 'check_webext_manifest__4b1e45b68c5d85573f69177601102dcb', 'check_eslint_extension__d09f7694', 'check_table_top_left__2d408821', 'check_extension_manifest__3d9fc0c33aef9d5f768043f561973d63', 'check_bookmark_folder_contains_url__3c0977b51a553edc4a7433583af5fcf6', 'is_expected_bookmarks__7a5a7856_aug_6', 'check_partial_extensions__2f01adaf', 'check_html_file_exists__48e4da460d6f3e132e4d3cc48ac9ca24', 'check_extension_name__a4b0d616259699774cd5ee2679c8a96c', 'check_gdrive_pdf_info__bb91e7693f2a30704f1d1cc79be73950', 'check_chrome_bookmark_exists__b5b6115e', 'check_gdrive_pdf_in_folder__661594c9', 'check_author_table_structure__f918ebce432b3a8956b1f7dd26a64d11', 'check_java_extension__c2202196', 'check_googledrive_files_exist__ba61b137046435f47239f8911466a875', 'check_extension_not_installed__cd5e025975a2ca9b4b85d9a92890bab8', 'check_last_table_structure__e0b81f05', 'is_extension_installed__f0c7e5c8', 'check_gdrive_pdf__a70aba9d', 'check_url_deleted__b835bc31', 'check_gdrive_pdf__bf7218b9', 'check_url_replaced_with_doi__e1878192', 'check_search_engine_changed__8a4fe4a4', 'check_table_height__60d42be3', 'check_gdrive_files__5876322dac1e30402171f6bcd2edb019', 'check_pdf_in_gdrive__a1937ab0', 'check_bookmark_folder_with_url__7a5a7856f1b642a4ade91ca81ca0f263000020251221151547', 'check_gitlens_extension__136a9bb5', 'check_vim_tabstop_enabled__76faceaeb09651393a582783b81e5798', 'check_gdrive_pdf_in_folder__ada1e84e', 'is_extension_installed__fdadaf5d', 'check_chrome_ext_single__abfa4b043befc7603aff32afef71bd1b', 'check_gdrive_files__12a331f2', 'check_table_top_position__5ff0bd58', 'check_html_files_exist__98346f3b', 'is_extension_installed__ba154228', 'check_csv_table_biz__fb9ae858', 'check_gdrive_pdf__82f685a6', 'check_chrome_experiments_exact_match__d7481e87a8a6dde50669ce517c215edf', 'check_table_count__4217a000', 'is_expected_bookmarks__2ad9387a_65d8_4e33_ad5b_7580065a27ca_task_verify_3', 'check_vim_hlsearch_enabled__5438ce42ea45fa77c023ecd730e398a5', 'check_liveserver_extension__a342948a', 'is_extension_installed__9333e0a8', 'check_table_row_added__fcf0fdbb', 'is_expected_search_query__f9cee6e4', 'is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_task_verify_4', 'check_gdrive_eml_files__74b11cf6', 'check_file_exists_webp__82cfb1cc', 'check_table_completeness__5067b5ea', 'check_webext_manifest__b9b28e4a', 'check_csv_table_mech__9629fd7e', 'check_recreation_html_element__dbcda1af4d2b810312479cfb54a15ab6', 'check_chrome_setting__df4ecdef', 'check_extension_uppercase__de1c34d5', 'check_extension_enabled__6f71517e0c42749af1a6363d9f36e224', 'check_recreation_html_element__fa1e76c31141f93d38de11c4bb8239cf', 'is_extension_installed__1d58e082', 'check_recreation_html_element__ed948b8fac72e2c98dd5028aabd38bf3', 'check_restclient_extension__d9713508', 'check_gdrive_pdf__f93e2dd2', 'check_table_sorted_with_row_integrity__b82f9c9e', 'check_table_row_count_by_index__f500041f', 'is_expected_bookmarks__ee2a139f7ed95e764d31e544c144187a', 'is_extension_installed__9024492c', 'check_chrome_experiments_contains_any__70b32da51f968d0aa45e836eecc52bdb', 'check_contains_url__a2f501d6', 'check_extension_version__a366045b', 'check_gdrive_file_exists__b791b56367b138c183fb013d8a0662b9', 'check_webext_manifest__81e9cd46', 'check_extension_name_path__d51a0e10', 'check_chrome_experiments_contains_all__e998f78abb27064086318477b860256b', 'is_expected_search_query__8164e914', 'check_gdrive_files__f478ac41', 'check_csv_table_bio__a4ea4d07', 'check_csv_table_cs3y__26538d5a', 'is_expected_bookmarks__2ad9387a_65d8_4e33_ad5b_7580065a27ca_aug_1_task_verify_1', 'check_html_contains__94d6a0b8', 'check_websites_valid__af875e9f', 'check_gdrive_files__7fbebc15', 'check_chrome_hardware_accel__8f79fa7c', 'check_csv_export_and_chrome__4e1aca7b', 'check_table_count_only__db5e4b5e', 'check_extension_description__db4e5321', 'check_table_centered__bbe9b961', 'check_account_server_url__9e272861', 'check_extension_version__407be0458b7b234fb5401d66a10f5221', 'check_csv_table_cs4y__5116177c', 'check_extension_enabled__ae6416e4', 'is_expected_bookmarks__a82b78bb_aug14', 'is_expected_bookmarks__35253b65', 'check_gdrive_file__6bd409d3', 'check_csv_file_and_tab__29a47154', 'check_docker_extension__acbefae0', 'check_extension_manifest__b5fa3477caae1bab3fd0d8b19ef4dfc7', 'check_url_contains__9fa2717e', 'is_expected_search_query__df8126bf', 'check_search_engine__d2ec4a7b', 'is_expected_url_pattern_match__907d006d1310a883a67bb7931a0dbae9', 'check_webext_manifest__c1aa21b2', 'check_extension_name__e5eabc9a', 'check_table_width__f4eb9543', 'check_all_urls_contain__82d38938', 'check_chrome_setting__8cfb2b13', 'check_tabs_and_bookmark__7a5a7856f1b642a4ade91ca81ca0f263000420251221151547', 'check_chrome_ext_productivity__3ce2c199effd9eefc89344b99f4cd5a7', 'check_googledrive_file_count_and_names__dd26081a', 'is_extension_installed__d374210b', 'check_webext_manifest__2ec2a9a8', 'check_url__5525d1c8', 'check_product_revenue_table__116fbd3460dada9be380bc664e224d99', 'check_gdrive_files__3dcb78c90ac64690f9f399090d59db08', 'is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_14_task_verify_1', 'check_chrome_tabs_and_bookmark__6d016e82', 'check_gdrive_files__fd7bb036', 'check_webext_manifest__c41914cdb60620ceca4e49be63e9e04c', 'check_gdrive_pdf_in_folder__9b0aebe9', 'check_grf_table__d24a6b5d9bcf09db9cd5a9cf07143d27', 'check_chrome_setting__1be3beb2', 'check_extension_manifest__986ae666fddb67409d6f0937ecd8b146', 'check_url_match__3b2adddad1034fb6d584d5542b1c7034', 'check_googledrive_folder__a18b8359', 'check_webext_manifest__60d86cd5', 'check_channel_revenue_table__22652fc14c0135d2f835816b613f7f8c', 'check_extension_enabled__1e533eef', 'check_file_executable__2bea57f7', 'check_chrome_setting__5936c775', 'check_csv_table_phys__5e1848a9', 'check_recreation_url__6dc3893f96ccd943c500af1756962de6', 'check_gdrive_pdf_in_folder__a29f27f0', 'check_prettier_extension__4d8ac023', 'check_gdrive_pdf__75da4280', 'check_chrome_setting__fd2ee811', 'check_webext_manifest__e32cecfb', 'check_gdrive_files__839d38d5', 'check_extension_version__f107c64c583fbfb012ea91827c8f61e3', 'check_extension_version__ae6416e4', 'check_extension_details__ae6416e4', 'check_extension_manifest_version__b7035c23581ef82d8725cee1b3aa987f', 'check_chrome_ext_last_two__f117f7ab324ab80f0f8ad254a42cb210', 'check_html_exists__7b7b0b2f', 'check_gdrive_pattern__6538c56492a239416edab32324471fe6', 'check_recreation_search__e0d56e1bf714f2c8f287d0502b986b63', 'check_table_bottom_left__6a4e1dd4', 'check_gdrive_pdf__c340b25e', 'check_cpp_extension__d45d3417', 'check_git_remote_url__f8582c17', 'check_amoxicillin_url__b070486d', 'check_jupyter_extension__bb12eb1b', 'is_expected_bookmarks__a82b78bb', 'check_gdrive_nested_files__688ade84b6705b21f2a27dc4863ba216', 'check_table_h_centered__dd5268aa', 'check_extension__e13be972', 'check_month_revenue_table__321ed1a008442c1379822ece597d5a12', 'check_table_row_count_by_index__83874f16', 'check_csv_table_med__1bf7a84e', 'check_extension_manifest_version__ae6416e4', 'check_chrome_setting__fca61186', 'check_html_conversion__b5ca523a54cc6ae837c08b60601febd6', 'is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_task_verify_1', 'check_extension_count__4329ad6c', 'check_chrome_setting__a785d845', 'check_csv_table_ee__c4163e13', 'check_websites_filled__53b95efcf5f95ba5e65301e716b45c01', 'check_extension_manifest__842d772f81b13ef554c1bda7e1c59bb7', 'check_chrome_extension_manifest__3d480472b35ce7003ca943ba6b2307fa', 'check_gdrive_pdf_in_folder__adf6c9b9', 'check_chrome_third_party_cookies__1394774d', 'check_gdrive_folder__55cd0f01', 'check_vim_tabstop__e2e649c5246b836f874ad28b723333bc', 'check_min_extension_count__0f84311e', 'check_extension_contains__fc6800da1dfd116cace0d10a635a3df0', 'check_extension_description__ae6416e4', 'check_table_completeness__5a197d93', 'check_webext_manifest__d1a4f844666cd6560872c801fdefe60a', 'check_at_least_one_extension__33e136ec', 'is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_14_task_verify_0', 'check_chrome_startup_urls__715fc21b1c707dd79fc5ab9b6e4df514', 'check_webext_manifest__70dff70b667529be05282f3babd2fe6e', 'check_gdrive_pdf_in_folder__42b1e128', 'is_extension_installed__2dcff8b4', 'check_gdrive_files_exist__74b11cf6', 'check_extension_source_type__ae6416e4', 'check_chrome_restore_setting__9cb5a6acff7768f53f96a10b1f86ac95', 'check_chrome_ext_subset__d7f8f8dc6935bc142e9b5dc629fdadfe', 'check_git_url_contains__07fbe1c4', 'check_table_left_aligned__b943f0fd', 'is_expected_bookmarks__2ad9387a', 'check_gdrive_files__59579682', 'check_chrome_extension_manifest__1e834a0e2fee4e317d31e5d8fca95c5c', 'check_table_right_aligned__d963ccf8', 'check_block_third_party_cookies__d0393d13df595b6db99860dc4f30ea7b', 'check_chrome_enhanced_safe_browsing__a6a55567']

def is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_5(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.
    This is a fixed version that uses the configurable 'names' parameter instead of hardcoded folder name.
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'liked_authors_websites_urls':
        expected_folder_names = rule.get('names', [])
        if not expected_folder_names:
            logger.error("Rule 'names' parameter is empty")
            return 0.0
        target_folder = None
        for bookmark in bookmarks['bookmark_bar']['children']:
            if bookmark['type'] == 'folder' and bookmark['name'] in expected_folder_names:
                target_folder = bookmark
                break
        if target_folder:
            logger.info(f"'{target_folder['name']}' folder exists")
            folder_urls = [bookmark['url'] for bookmark in target_folder['children'] if bookmark['type'] == 'url']
            logger.info(f"Here is the '{target_folder['name']}' folder's urls: {folder_urls}")
            urls = rule['urls']
            for (idx, url) in enumerate(urls):
                if isinstance(url, str):
                    urls[idx] = [url]
            combinations = product(*urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    return 1.0
            return 0.0
        else:
            logger.info(f'Expected folder not found. Looking for one of: {expected_folder_names}')
            return 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_webext_manifest__bba937aa(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    expected_has_blank_description = expected.get('has_blank_description', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    if expected_has_blank_description:
        checks += 1
        description = manifest.get('description', '')
        if not description or not description.strip():
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_chrome_experiments_not_contains__f72b4886483f6a090bcc6a28ba49acde(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the enabled Chrome experiments do NOT contain the expected experiment.

    Args:
        result: List of enabled experiment names
        expected: Dict with 'experiment_name' key specifying the experiment that should NOT be enabled
        **options: Additional options

    Returns:
        1.0 if the expected experiment is NOT in the list, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    forbidden_experiment = expected.get('experiment_name', '')
    if not forbidden_experiment:
        logger.error('No experiment_name specified in expected rules')
        return 0.0
    if forbidden_experiment not in result:
        logger.info(f"Confirmed experiment '{forbidden_experiment}' is NOT in enabled experiments")
        return 1.0
    else:
        logger.info(f"Experiment '{forbidden_experiment}' found but should be disabled. Enabled: {result}")
        return 0.0

def check_gdrive_files__da0d777d0a1e3c8d735dbb81a2e45a5c(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected files exist in Google Drive folder.

    Args:
        result: List of filenames found in Google Drive folder
        expected: Dict with 'expected_files' key containing list of required filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.warning('No expected files specified')
        return 0.0
    if not result:
        logger.warning('No files found in Google Drive')
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
            logger.info(f'Found expected file: {expected_file}')
        else:
            logger.warning(f'Missing expected file: {expected_file}')
    score = found_count / len(expected_files)
    logger.info(f'Found {found_count}/{len(expected_files)} expected files. Score: {score}')
    return score

def check_gdrive_pdf_in_folder__b2131ee6(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_gdrive_pdf_in_folder__94eb0e23(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_exact_bookmarks__b53050a013ee52945e88dd9cbd9764a5(bookmarks: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Checks that required URLs are present and forbidden URLs are absent in bookmark bar.

    Args:
        bookmarks: Bookmarks data structure from get_bookmarks getter
        expected: Expected configuration with:
            - required_urls (List[str]): URLs that must be in bookmark bar
            - forbidden_urls (List[str]): URLs that must NOT be in bookmark bar
        **options: Additional options (unused)

    Returns:
        float: 0.5 if required URLs present, 1.0 if required URLs present and forbidden URLs absent, 0.0 otherwise
    """
    if not bookmarks:
        logger.info('No bookmarks data available')
        return 0.0
    required_urls = expected.get('required_urls', [])
    forbidden_urls = expected.get('forbidden_urls', [])
    bookmark_bar = bookmarks.get('bookmark_bar', {})
    children = bookmark_bar.get('children', [])
    bookmark_bar_urls = [child.get('url') for child in children if child.get('type') == 'url']
    logger.info(f'URLs in bookmark bar: {bookmark_bar_urls}')
    logger.info(f'Required URLs: {required_urls}')
    logger.info(f'Forbidden URLs: {forbidden_urls}')
    all_required_present = all((url in bookmark_bar_urls for url in required_urls))
    if not all_required_present:
        logger.info('Not all required URLs are present in bookmark bar')
        return 0.0
    logger.info('All required URLs are present in bookmark bar')
    any_forbidden_present = any((url in bookmark_bar_urls for url in forbidden_urls))
    if any_forbidden_present:
        logger.info('Some forbidden URLs are present in bookmark bar')
        return 0.5
    else:
        logger.info('No forbidden URLs are present in bookmark bar')
        return 1.0

def check_table_count_and_content__638d8a63(result, expected, **options):
    """Check if document has correct table count and first table content.

    Args:
        result: Dict from getter with 'table_count' and 'first_table_data'
        expected: Dict with expected 'table_count' and 'first_table_data'
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.error('Invalid result or expected format')
        return 0.0
    score = 0.0
    result_count = result.get('table_count', 0)
    expected_count = expected.get('table_count', 0)
    if result_count == expected_count:
        score += 0.5
        logger.info(f'✓ Table count matches: {result_count}')
    else:
        logger.info(f'✗ Table count mismatch: got {result_count}, expected {expected_count}')
    result_data = result.get('first_table_data', [])
    expected_data = expected.get('first_table_data', [])
    if len(result_data) == len(expected_data):
        matching_rows = 0
        for (i, (result_row, expected_row)) in enumerate(zip(result_data, expected_data)):
            if len(result_row) == len(expected_row):
                if all((r.strip().lower() == e.strip().lower() for (r, e) in zip(result_row, expected_row))):
                    matching_rows += 1
        if matching_rows == len(expected_data):
            score += 0.5
            logger.info(f'✓ First table content matches')
        else:
            partial = matching_rows / len(expected_data) * 0.5 if expected_data else 0
            score += partial
            logger.info(f'Partial match: {matching_rows}/{len(expected_data)} rows correct')
    else:
        logger.info(f'✗ First table row count mismatch: got {len(result_data)}, expected {len(expected_data)}')
    return score

def check_webext_manifest__dea4a5ea3d588b414bf151fee72b35b0(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if web extension manifest.json matches expected structure.

    Args:
        result: Manifest JSON dict from getter (or None if file not found)
        expected: Expected manifest structure with rules to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.warning('Manifest file not found')
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_description = expected.get('has_description')
    expected_has_content_scripts = expected.get('has_content_scripts')
    score = 0.0
    total_checks = 0
    if expected_name is not None:
        total_checks += 1
        if result.get('name') == expected_name:
            score += 1
            logger.info(f'Name matches: {expected_name}')
        else:
            logger.warning(f"Name mismatch: expected '{expected_name}', got '{result.get('name')}'")
    if expected_version is not None:
        total_checks += 1
        if result.get('version') == expected_version:
            score += 1
            logger.info(f'Version matches: {expected_version}')
        else:
            logger.warning(f"Version mismatch: expected '{expected_version}', got '{result.get('version')}'")
    if expected_has_background:
        total_checks += 1
        if 'background' in result and 'scripts' in result['background'] and (len(result['background']['scripts']) > 0):
            score += 1
            logger.info(f"Background scripts found: {result['background']['scripts']}")
        else:
            logger.warning('Background scripts not found')
    if expected_has_browser_action:
        total_checks += 1
        if 'browser_action' in result:
            score += 1
            logger.info('Browser action found')
        else:
            logger.warning('Browser action not found')
    if expected_has_page_action:
        total_checks += 1
        if 'page_action' in result:
            score += 1
            logger.info('Page action found')
        else:
            logger.warning('Page action not found')
    if expected_has_description is not None:
        total_checks += 1
        description = result.get('description', '')
        if expected_has_description is False:
            if not description or description.strip() == '':
                score += 1
                logger.info('Description is absent or empty as expected')
            else:
                logger.warning(f"Description should be absent but found: '{description}'")
        elif description and description.strip() != '':
            score += 1
            logger.info(f"Description found: '{description}'")
        else:
            logger.warning('Description is missing but was expected')
    if expected_has_content_scripts is not None:
        total_checks += 1
        if expected_has_content_scripts is False:
            if 'content_scripts' not in result:
                score += 1
                logger.info('Content scripts are absent as expected')
            else:
                logger.warning(f"Content scripts should be absent but found: {result['content_scripts']}")
        elif 'content_scripts' in result:
            score += 1
            logger.info(f"Content scripts found: {result['content_scripts']}")
        else:
            logger.warning('Content scripts are missing but were expected')
    if total_checks == 0:
        logger.warning('No checks to perform')
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks} checks passed)')
    return final_score

def check_extension_manifest__1edda3bd4444a8fb1554a227ff178017(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate browser extension manifest.json structure.

    Args:
        result: Manifest data from getter
        expected: Rules dict containing expected values for name, version, background, browser_action
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct field)
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dictionary')
        return 0.0
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    check_background = expected.get('check_background', False)
    check_browser_action = expected.get('check_browser_action', False)
    score = 0.0
    total_checks = 0
    total_checks += 1
    if result.get('name') == expected_name:
        score += 0.2
        logger.info(f'Name matches: {expected_name}')
    else:
        logger.info(f"Name mismatch: expected={expected_name}, got={result.get('name')}")
    total_checks += 1
    if result.get('version') == expected_version:
        score += 0.2
        logger.info(f'Version matches: {expected_version}')
    else:
        logger.info(f"Version mismatch: expected={expected_version}, got={result.get('version')}")
    total_checks += 1
    description = result.get('description', '')
    if description == '' or description is None:
        score += 0.2
        logger.info('Description is blank/empty as required')
    else:
        logger.info(f'Description should be blank but got: {description}')
    if check_background:
        total_checks += 1
        background = result.get('background', {})
        if isinstance(background, dict) and 'scripts' in background:
            scripts = background.get('scripts', [])
            if isinstance(scripts, list) and len(scripts) > 0:
                score += 0.2
                logger.info(f'Background scripts found: {scripts}')
            else:
                logger.info('Background scripts field exists but is empty')
        else:
            logger.info('Background section missing or invalid')
    if check_browser_action:
        total_checks += 1
        browser_action = result.get('browser_action', {})
        if isinstance(browser_action, dict) and len(browser_action) > 0:
            has_fields = any((k in browser_action for k in ['default_popup', 'default_icon', 'default_title']))
            if has_fields:
                score += 0.2
                logger.info(f'Browser action found with fields: {list(browser_action.keys())}')
            else:
                logger.info('Browser action exists but missing expected fields')
        else:
            logger.info('Browser action section missing or invalid')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_gdrive_pdf__3a6a494e(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_gdrive_pdf__9d82c3b2(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def is_expected_url_pattern_match__ac7e6c6cfb1c60f162b6d2bc92009d82(result, expected, **options) -> float:
    """
    Checks if the active tab URL matches multiple regex patterns.

    Args:
        result: The active tab info (string URL or dict with 'url' field)
        expected: Dictionary with 'expected' key containing list of regex patterns
        **options: Additional options

    Returns:
        float: 1.0 if all patterns match, 0.0 otherwise
    """
    if not result:
        logger.info('[PATTERN_MATCH] No result provided')
        return 0.0
    if isinstance(result, str):
        result_url = result
    elif isinstance(result, dict) and 'url' in result:
        result_url = result['url']
    else:
        logger.error(f'[PATTERN_MATCH] Invalid result format: {type(result)}')
        return 0.0
    logger.info(f'[PATTERN_MATCH] Result URL: {result_url}')
    patterns = expected.get('expected', [])
    logger.info(f'[PATTERN_MATCH] Expected patterns: {patterns}')
    for pattern in patterns:
        match = re.search(pattern, result_url)
        if not match:
            logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' not found in URL")
            return 0.0
        logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' matched")
    logger.info('[PATTERN_MATCH] All patterns matched successfully')
    return 1.0

def check_gdrive_files__45753efe(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_chrome_extension_manifest__6e0f538e8ca610f13ae0673161924889(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate Chrome extension manifest and files.

    Args:
        result: Dict from getter with 'manifest', 'files_exist', 'all_files_exist'
        expected: Dict with expected manifest structure rules:
            - manifest_checks: List of dicts with 'key' (list of nested keys) and 'value' (expected value)
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    manifest = result.get('manifest', {})
    all_files_exist = result.get('all_files_exist', False)
    file_score = 0.5 if all_files_exist else 0.0
    manifest_checks = expected.get('manifest_checks', [])
    if not manifest_checks:
        return file_score
    passed_checks = 0
    total_checks = len(manifest_checks)
    for check in manifest_checks:
        key_path = check.get('key', [])
        expected_value = check.get('value')
        current = manifest
        try:
            for key in key_path:
                current = current[key]
            if current == expected_value:
                passed_checks += 1
        except (KeyError, TypeError):
            pass
    manifest_score = 0.5 * (passed_checks / total_checks) if total_checks > 0 else 0.0
    return file_score + manifest_score

def is_expected_url_pattern_match_or__b070486d(result, rules) -> float:
    """
    Check if the result URL matches ANY of the expected regex patterns (OR logic).

    This function is used to search for expected patterns in the URL using regex.
    result is the return value of function "get_active_url_from_accessTree" or similar.

    Args:
        result: Either a string URL or a dict with 'url' field
        rules: Dict containing 'expected' key with list of regex patterns

    Returns:
        float: 1.0 if ANY pattern matches, 0.0 otherwise
    """
    if not result:
        logger.info('Result is empty')
        return 0.0
    if isinstance(result, str):
        result_url = result
        logger.info(f'Result URL (string): {result_url}')
    elif isinstance(result, dict) and 'url' in result:
        result_url = result['url']
        logger.info(f'Result URL (dict): {result_url}')
    else:
        logger.error(f"Invalid result format: {type(result)}, expected string URL or dict with 'url' field")
        return 0.0
    patterns = rules.get('expected', [])
    if not patterns:
        logger.error('No expected patterns provided in rules')
        return 0.0
    logger.info(f'Expected patterns: {patterns}')
    for pattern in patterns:
        try:
            match = re.search(pattern, result_url)
            if match:
                logger.info(f"Pattern '{pattern}' matched URL '{result_url}'")
                return 1.0
            else:
                logger.debug(f"Pattern '{pattern}' did not match URL '{result_url}'")
        except re.error as e:
            logger.error(f"Invalid regex pattern '{pattern}': {e}")
            continue
    logger.info(f"No patterns matched URL '{result_url}'")
    return 0.0

def check_chrome_ext_middle__484fc06534818b90d35ad638d0805c7c(result: List[str], expected: Dict, **options) -> float:
    """Check if expected Chrome extensions are installed.

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' key containing list of required extension names
        **options: Additional options (not used)

    Returns:
        1.0 if all expected extensions are installed, 0.0 otherwise
    """
    if not result:
        logger.info('No extensions found in result')
        return 0.0
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        logger.warning('No expected extensions specified')
        return 0.0
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Expected extensions: {expected_extensions}')
    set_expected = set(expected_extensions)
    set_installed = set(result)
    if set_expected.issubset(set_installed):
        logger.info('All expected extensions are installed')
        return 1.0
    else:
        missing = set_expected - set_installed
        logger.info(f'Missing extensions: {missing}')
        return 0.0

def check_extension_count__7a8c3f12(result: str, expected: dict, **options) -> float:
    """Check if the number of installed extensions meets minimum requirement.

    Args:
        result: Output from 'code --list-extensions' command (one extension per line)
        expected: Dict with 'min_count' key specifying minimum number of extensions

    Returns:
        1.0 if extension count >= min_count, 0.0 otherwise
    """
    if not result:
        return 0.0
    min_count = expected.get('min_count', 1)
    lines = [line.strip() for line in result.strip().split('\n') if line.strip()]
    actual_count = len(lines)
    return 1.0 if actual_count >= min_count else 0.0

def check_table_count_only__d0985f12(result, expected, **options):
    """Check table count and verify merged vowel table content."""
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.error('Result or expected is not a dict')
        return 0.0
    result_count = result.get('table_count', 0)
    expected_count = expected.get('table_count', 0)
    has_merged_table = result.get('has_merged_table', False)
    merged_table_vowels = result.get('merged_table_vowels', [])
    merged_table_rows = result.get('merged_table_rows', 0)
    original_vowel_tables_found = result.get('original_vowel_tables_found', 0)
    total_original_rows = result.get('total_original_rows', 0)
    merged_table_position = result.get('merged_table_position', -1)
    first_vowel_table_position = result.get('first_vowel_table_position', -1)
    content_verification_passed = result.get('content_verification_passed', False)
    expected_vowels = expected.get('vowel_patterns', ['a', 'e', 'i', 'o', 'u'])
    min_expected_rows = expected.get('min_merged_rows', 4)
    score = 0.0
    max_score = 1.0
    if not has_merged_table:
        logger.info(f'✗ FAIL: No merged table found - cannot proceed with verification')
        return 0.0
    if original_vowel_tables_found >= 4:
        logger.info(f'✗ FAIL: All original vowel tables still exist ({original_vowel_tables_found} found) - merge did not happen')
        return 0.0
    if result_count == expected_count:
        logger.info(f'✓ Table count matches: {result_count} (expected {expected_count})')
        score += 0.2
    else:
        logger.info(f'✗ Table count mismatch: got {result_count}, expected {expected_count}')
    if set(merged_table_vowels) >= set(expected_vowels):
        logger.info(f'✓ Merged table contains all vowel patterns: {merged_table_vowels}')
        score += 0.25
    else:
        missing = set(expected_vowels) - set(merged_table_vowels)
        logger.info(f'✗ FAIL: Merged table missing vowel patterns: {missing}')
        return 0.0
    if first_vowel_table_position >= 0:
        position_diff = abs(merged_table_position - first_vowel_table_position)
        if position_diff <= 3:
            logger.info(f'✓ Merged table is at expected position (diff: {position_diff})')
            score += 0.15
        else:
            logger.info(f'⚠ Merged table position ({merged_table_position}) is far from first vowel table ({first_vowel_table_position}), diff: {position_diff}')
            if position_diff <= 5:
                score += 0.1
            else:
                score += 0.05
    else:
        logger.info(f'⚠ Could not determine expected position')
        score += 0.05
    if content_verification_passed:
        logger.info(f'✓ Content from original tables verified in merged table')
        score += 0.2
    else:
        logger.info(f"✗ FAIL: Content verification failed - merged table doesn't contain content from original tables")
        return 0.0
    if total_original_rows > 0:
        tolerance = 0.35
        lower_bound = total_original_rows * (1 - tolerance)
        upper_bound = total_original_rows * (1 + tolerance)
        if lower_bound <= merged_table_rows <= upper_bound:
            logger.info(f'✓ Merged table row count ({merged_table_rows}) is within expected range [{lower_bound:.1f}, {upper_bound:.1f}] based on original tables ({total_original_rows} total rows)')
            score += 0.2
        elif merged_table_rows >= min_expected_rows:
            logger.info(f'⚠ Merged table row count ({merged_table_rows}) is outside expected range [{lower_bound:.1f}, {upper_bound:.1f}], but meets minimum ({min_expected_rows})')
            score += 0.1
        else:
            logger.info(f'✗ FAIL: Merged table rows ({merged_table_rows}) insufficient, expected ~{total_original_rows} rows')
            return 0.0
    elif merged_table_rows >= min_expected_rows:
        logger.info(f'✓ Merged table has sufficient rows: {merged_table_rows} (min {min_expected_rows})')
        score += 0.2
    else:
        logger.info(f'✗ FAIL: Merged table rows insufficient: {merged_table_rows} < {min_expected_rows}')
        return 0.0
    if original_vowel_tables_found == 0:
        logger.info(f'✓ All original vowel tables were properly merged (0 remaining)')
    elif original_vowel_tables_found <= 1:
        logger.info(f'⚠ Most vowel tables merged ({original_vowel_tables_found} remaining)')
        score *= 0.95
    elif original_vowel_tables_found <= 2:
        logger.info(f'⚠ Some original vowel tables remaining ({original_vowel_tables_found} found)')
        score *= 0.85
    else:
        logger.info(f'⚠ Many original vowel tables remaining ({original_vowel_tables_found} found)')
        score *= 0.7
    logger.info(f'Final score: {score:.2f}/{max_score}')
    return score

def check_extension_description__221bd7ef80e48e21d62e7d38635cf26c(result, expected, **options):
    """Check if the extension description matches the expected description.

    Args:
        result: Description string from getter
        expected: Dict with 'description' key containing the expected description
        **options: Additional options (not used)

    Returns:
        float: 1.0 if descriptions match, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_description = expected.get('description', '')
    if not expected_description:
        return 0.0
    if result == expected_description:
        return 1.0
    else:
        return 0.0

def check_gdrive_pdf_in_folder__41049d6b(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def is_expected_search_query__e346411e(active_tab_info: Dict[str, str], rules: Dict[str, Any]) -> float:
    """
    Check if the active tab URL matches the expected search query pattern.
    Variation 8: Search for first name from cell B3 (Mara).
    """
    if not active_tab_info:
        return 0.0
    expected = rules['expect']
    pattern = expected['pattern']
    matched = re.search(pattern, active_tab_info['url'])
    if matched:
        return 1.0
    return 0.0

def check_gdrive_pdf_in_folder__db0d2d11(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_csv_table_civil__2c0c636e(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_cookies_deleted__4730218c(cookie_data, expected):
    """
    Check if cookies have been deleted from Chrome.

    This function verifies that the cookie count is at or below a specified threshold.
    When type='rule', the framework passes the rules dict directly as expected parameter.

    Args:
        cookie_data: List of cookies returned from get_cookie_data getter, or None if DB doesn't exist
        expected: Dict containing the rules directly (when type='rule'):
                  {
                      "max_cookie_count": <int>  # Maximum allowed cookies (0 means all deleted)
                  }

    Returns:
        float: 1.0 if cookie count <= max_cookie_count, 0.0 otherwise
    """
    if cookie_data is None:
        return 1.0
    max_cookie_count = expected.get('max_cookie_count', 0)
    actual_cookie_count = len(cookie_data)
    if actual_cookie_count <= max_cookie_count:
        return 1.0
    return 0.0

def is_extension_installed__72cc5a03(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_extension_installed__c548e99163664e4ea2d87d3cc236c650(actual: str, expected: Dict, **options) -> float:
    """
    Check if the VSCode extension is installed.

    Args:
        actual (str): Output from 'code --list-extensions' command
        expected (Dict): Expected rules with 'extension_id' field
        **options: Additional options

    Returns:
        float: 1.0 if extension is installed, 0.0 otherwise
    """
    if not actual:
        return 0.0
    extension_id = expected.get('extension_id', '')
    if not extension_id:
        return 0.0
    installed_extensions = [line.strip() for line in actual.strip().split('\n') if line.strip()]
    for ext in installed_extensions:
        if ext.lower() == extension_id.lower():
            return 1.0
    return 0.0

def check_chrome_setting__eec6beed(result, expected, **options):
    """
    Check if the Chrome Translation settings setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_html_conversion_and_tab__6920e35a(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if HTML file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with html_path and html_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f'HTML file exists: +0.5')
    else:
        logger.warning(f'HTML file does not exist')
    chrome_tabs = result.get('chrome_tabs', [])
    html_url_pattern = expected.get('html_url_pattern', '')
    html_opened = False
    for tab_url in chrome_tabs:
        if html_url_pattern in tab_url:
            html_opened = True
            break
    if html_opened:
        score += 0.5
        logger.info(f'HTML file opened in Chrome: +0.5')
    else:
        logger.warning(f'HTML file not found in Chrome tabs. Tabs: {chrome_tabs}')
    logger.info(f'Final score: {score}')
    return score

def check_gdrive_pdf__9e5c8fdc(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_chrome_setting__d84aae11(result, expected, **options):
    """
    Check if the Chrome Startup behavior setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_startup_new_tab__bd145bbb387b98da7ba508028ec58276(result, expected, **options):
    """
    Check if Chrome is set to open new tab page on startup.

    Args:
        result: String returned from get_new_startup_page getter ("true" or "false")
        expected: Expected value from rules (dict with "expected" key)
        **options: Additional options

    Returns:
        float: 1.0 if setting matches expected, 0.0 otherwise
    """
    expected_value = expected.get('expected', 'true')
    if result == expected_value:
        return 1.0
    else:
        return 0.0

def check_recreation_html_element__2aa2c046ed13ceae9537c9a8e566d558(result, expected, **options):
    """
    Verify that the task to find the nearest open slot for Antelope Island was completed.
    Checks both location correctness and availability data, not just generic page structure.

    Args:
        result: Dict from getter function containing:
            - location_verified: bool (True if Antelope Island was searched)
            - availability_found: bool (True if availability table was found)
            - url: str (current page URL)
            - search_term: str or None (extracted location name)
            - first_available_date: str or None (first available slot)
        expected: Expected dict structure from rules
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    logger.info(f'[DEBUG] check_recreation_html_element called with result: {result}')
    logger.info(f'[DEBUG] check_recreation_html_element called with expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, got {type(result)}: {result}')
        return 0.0
    score = 0.0
    max_score = 0.0
    max_score += 0.4
    location_verified = result.get('location_verified', False)
    if location_verified:
        score += 0.4
        logger.info('[DEBUG] Location verification PASSED: Antelope Island was searched')
    else:
        logger.info('[DEBUG] Location verification FAILED: Antelope Island not found in URL or page content')
    max_score += 0.4
    availability_found = result.get('availability_found', False)
    if availability_found:
        score += 0.4
        logger.info('[DEBUG] Availability check PASSED: Results table found')
    else:
        logger.info('[DEBUG] Availability check FAILED: No availability table found')
    max_score += 0.2
    first_available_date = result.get('first_available_date')
    if first_available_date:
        score += 0.2
        logger.info(f"[DEBUG] Available date check PASSED: Found date '{first_available_date}'")
    else:
        logger.info('[DEBUG] Available date check FAILED: No available date found (may still be acceptable)')
    url = result.get('url', 'N/A')
    search_term = result.get('search_term', 'N/A')
    logger.info(f'[DEBUG] URL: {url}')
    logger.info(f'[DEBUG] Search term: {search_term}')
    logger.info(f'[DEBUG] Final score: {score}/{max_score}')
    if score >= 0.8:
        logger.info('[DEBUG] Task COMPLETED: Nearest open slot for Antelope Island was found')
        return 1.0
    else:
        logger.info(f'[DEBUG] Task INCOMPLETE: Score {score} < 0.8 threshold')
        return score

def check_vim_hlsearch__47d76eea6c988a4beb0cae6b4109cc24(result: str, rules: Dict[str, List[str]]) -> float:
    """
    Check if the result contains expected include strings and excludes unwanted strings.

    Args:
        result: Output from the getter function
        rules: Dictionary with 'include' and 'exclude' lists

    Returns:
        1.0 if all include strings are present and exclude strings are absent, 0.0 otherwise
    """
    if result is None:
        return 0.0
    logger.info(f'Checking result: {result}, rules: {rules}')
    include = rules.get('include', [])
    exclude = rules.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_recreation_html_element__6adae5a4637b133a46055994fbaa8dd4(result, expected, **options):
    """
    Verify that Bear Lake search was performed and available reservation results are displayed.

    This metric checks:
    1. Bear Lake was searched for (appears in page content, URL, or search inputs)
    2. Reservation results are displayed on the page
    3. Availability information is present (indicating available reservations exist)

    Args:
        result: Dict from getter function containing:
            - bear_lake_found: bool (Bear Lake appears in page)
            - has_results: bool (reservation results displayed)
            - has_availability: bool (availability info present)
            - reservation_dates: list (extracted dates, for additional verification)
            - url_contains_search: bool (URL has search params)
        expected: Expected dict structure from rules (all should be true for task completion)
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    logger.info(f'[DEBUG] check_recreation_html_element called with result: {result}')
    logger.info(f'[DEBUG] check_recreation_html_element called with expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, got {type(result)}: {result}')
        return 0.0
    score = 0.0
    max_score = 3.0
    if result.get('bear_lake_found', False):
        score += 1.0
        logger.info('[DEBUG] ✓ Bear Lake search verified')
    else:
        logger.info('[DEBUG] ✗ Bear Lake NOT found in page - search may not have been performed')
        return 0.0
    if result.get('has_results', False):
        score += 1.0
        logger.info('[DEBUG] ✓ Reservation results displayed')
    else:
        logger.info('[DEBUG] ✗ No reservation results found - search may have failed or returned no results')
        return 0.0
    if result.get('has_availability', False):
        score += 1.0
        logger.info('[DEBUG] ✓ Availability information present')
    else:
        logger.info('[DEBUG] ⚠ No availability information found - may indicate no reservations available')
    if result.get('url_contains_search', False):
        logger.info('[DEBUG] ✓ URL contains search parameters (bonus signal)')
    if result.get('reservation_dates') and len(result['reservation_dates']) > 0:
        logger.info(f"[DEBUG] ✓ Found {len(result['reservation_dates'])} date elements (bonus signal)")
    final_score = score / max_score
    logger.info(f'[DEBUG] Final score: {final_score:.2f} ({score}/{max_score} checks passed)')
    return final_score

def check_webext_manifest__88e21052(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    expected_description_empty = expected.get('description_empty', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    if expected_description_empty:
        checks += 1
        description = manifest.get('description', '')
        if description == '':
            score += 1.0
    return score / checks if checks > 0 else 0.0

def is_expected_search_query__cbf92d65(active_tab_info: Dict[str, str], rules: Dict[str, Any]) -> float:
    """
    Check if the active tab URL matches the expected search query pattern.
    Variation 4: Search for date from cell G6 (16/08/2016).
    """
    if not active_tab_info:
        return 0.0
    expected = rules['expect']
    pattern = expected['pattern']
    matched = re.search(pattern, active_tab_info['url'])
    if matched:
        return 1.0
    return 0.0

def check_backup_extension_files__8ece34a1(result, expected, **options):
    """Check if backup extension files match expected list.

    Args:
        result: List of filenames from getter
        expected: Dict with 'expected_files' key

    Returns:
        float: 1.0 if files match, 0.0 otherwise
    """
    expected_files = sorted(expected.get('expected_files', []))
    result_sorted = sorted(result)
    if result_sorted == expected_files:
        return 1.0
    else:
        print(f'Expected: {expected_files}')
        print(f'Got: {result_sorted}')
        return 0.0

def check_chrome_do_not_track__a858c66a(result, expected, **options):
    """Check if Chrome Do Not Track is in expected state.

    Args:
        result: Do Not Track settings dict from getter
        expected: Expected rules dict with 'enabled' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_enabled = expected.get('enabled', False)
    actual_enabled = result.get('enabled', False)
    logger.info(f'Expected Do Not Track enabled: {expected_enabled}')
    logger.info(f'Actual Do Not Track enabled: {actual_enabled}')
    return 1.0 if actual_enabled == expected_enabled else 0.0

def check_go_extension__303c34ab(actual: str, rules: Dict, **options):
    """
    Check if Go extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_table_top_right__199b082c(result, expected, **options):
    """Check if table is positioned in top-right corner.

    Args:
        result: Table position dict from getter
        expected: Dict with 'max_top' and 'min_left'
        **options: Additional options

    Returns:
        float: 1.0 if both conditions met, 0.0 otherwise
    """
    if result is None:
        return 0.0
    max_top = expected.get('max_top', 0)
    min_left = expected.get('min_left', 0)
    if result['top'] < max_top and result['left'] > min_left:
        return 1.0
    else:
        return 0.0

def check_extension_exact_match__6410307fcb9f7a3e584e1435af0837eb(actual: str, expected: dict, **options) -> float:
    """
    Check if the VSCode extension with exact ID format is installed.

    Args:
        actual (str): Output from 'code --list-extensions' command
        expected (dict): Expected rules with 'extension_id' field
        **options: Additional options

    Returns:
        float: 1.0 if extension ID is found exactly, 0.0 otherwise
    """
    if not actual:
        return 0.0
    extension_id = expected.get('extension_id', '')
    if not extension_id:
        return 0.0
    installed_extensions = [line.strip() for line in actual.strip().split('\n') if line.strip()]
    for ext in installed_extensions:
        if ext.lower() == extension_id.lower():
            return 1.0
    return 0.0

def check_extension_description__2224641eff529090c9cc8ad62127e176(result, expected, **options):
    """
    Check if the extension description matches the expected value.

    This metric verifies that:
    1. The extension was found and loaded by Chrome (result is not empty)
    2. The extension is enabled (verified by the getter)
    3. The description matches the expected value

    Args:
        result: String from getter containing the extension description
        expected: Dict with 'extension_description' key
        **options: Additional options (unused)

    Returns:
        float: 1.0 if descriptions match and extension is enabled, 0.0 otherwise
    """
    expected_description = expected.get('extension_description', '')
    if not result:
        return 0.0
    if result == expected_description:
        return 1.0
    return 0.0

def check_urls_present__92bc11a1(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if the result contains valid-looking URLs.

    Args:
        result: List of values from a column
        expected: Dict with 'min_count' and optional 'url_pattern' keys
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    min_count = expected.get('min_count', 1)
    url_pattern = expected.get('url_pattern', 'https?://|www\\.')
    url_count = 0
    for value in result:
        if re.search(url_pattern, value, re.IGNORECASE):
            url_count += 1
    if url_count >= min_count:
        return 1.0
    else:
        return url_count / min_count if min_count > 0 else 0.0

def is_expected_search_query__e154bb96(combined_state: Dict[str, Any], rules: Dict[str, Any]) -> float:
    """
    Check if cell E6 was read and Chrome search URL matches the expected query.
    This verifies two things:
    1. Cell E6 in the spreadsheet contains 'United States'
    2. Chrome active tab URL contains the search query for 'united states'

    Args:
        combined_state: Dictionary with 'cell_value' and 'active_url' keys
        rules: Dictionary with 'expect' containing 'cell_value' and 'pattern' keys

    Returns:
        float: 1.0 if both checks pass, 0.0 otherwise
    """
    if not combined_state:
        logger.warning('[SEARCH_QUERY_CHECK] Combined state is None or empty')
        return 0.0
    cell_value = combined_state.get('cell_value')
    active_url = combined_state.get('active_url')
    expected_cell_value = rules['expect'].get('cell_value', 'United States')
    cell_score = check_cell_value_matches_expected(cell_value, expected_cell_value)
    url_score = check_url_matches_pattern(active_url, rules)
    if cell_score == 1.0 and url_score == 1.0:
        logger.info('[SEARCH_QUERY_CHECK] Both cell value and URL checks passed')
        return 1.0
    logger.warning(f'[SEARCH_QUERY_CHECK] Checks failed - cell_score: {cell_score}, url_score: {url_score}')
    return 0.0

def check_chrome_experiments_contains__0a40109ea287ab7bbd8cd9175a7ce6a5(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the enabled Chrome experiments contain the expected experiment.

    Args:
        result: List of enabled experiment names
        expected: Dict with 'experiment_name' key specifying the expected experiment
        **options: Additional options

    Returns:
        1.0 if the expected experiment is in the list, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    expected_experiment = expected.get('experiment_name', '')
    if not expected_experiment:
        logger.error('No experiment_name specified in expected rules')
        return 0.0
    if expected_experiment in result:
        logger.info(f"Found expected experiment '{expected_experiment}' in enabled experiments")
        return 1.0
    else:
        logger.info(f"Expected experiment '{expected_experiment}' not found. Enabled: {result}")
        return 0.0

def check_bookmarks_created__5641e4c0(result, expected, **options):
    """Check if expected bookmarks were created.

    Args:
        result: List of bookmark dicts with 'title' and 'url' keys from getter
        expected: Dict with:
            - expected_count: Expected number of bookmarks
            - expected_urls: List of expected URLs

    Returns:
        Score between 0.0 and 1.0
    """
    expected_count = expected.get('expected_count', 2)
    expected_urls = expected.get('expected_urls', [])
    if not isinstance(result, list):
        return 0.0
    if len(result) != expected_count:
        return 0.0
    if expected_urls:
        result_urls = [bookmark.get('url', '') for bookmark in result if isinstance(bookmark, dict)]
        matches = sum((1 for url in expected_urls if url in result_urls))
        if matches == len(expected_urls):
            return 1.0
        else:
            return matches / len(expected_urls)
    else:
        return 1.0

def check_pdf_and_tabs__a852a89d(result, expected, **options):
    """Check if both PDF files exist and tab count matches.

    Args:
        result: Dict with 'pdf_count' and 'tab_count' keys
        expected: Dict with 'pdf_count' and 'tab_count' keys

    Returns:
        float: Score from 0.0 to 1.0
            - 0.5 for correct tab count
            - 0.5 for correct PDF count
    """
    expected_pdf_count = expected.get('pdf_count', 2)
    expected_tab_count = expected.get('tab_count', 2)
    score = 0.0
    if result.get('tab_count', 0) == expected_tab_count:
        score += 0.5
        logger.info(f"Tab count matches: {result.get('tab_count', 0)} == {expected_tab_count}")
    else:
        logger.warning(f"Tab count mismatch: {result.get('tab_count', 0)} != {expected_tab_count}")
    if result.get('pdf_count', 0) == expected_pdf_count:
        score += 0.5
        logger.info(f"PDF count matches: {result.get('pdf_count', 0)} == {expected_pdf_count}")
    else:
        logger.warning(f"PDF count mismatch: {result.get('pdf_count', 0)} != {expected_pdf_count}")
    return score

def check_csv_table_math__bdcc0312(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def is_extension_installed__7fabdee6(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_extension_count__0fb1bf36ddef8ab5a554de56ff7f0c3d(result, expected, **options):
    """Check if the number of unpacked extensions meets the minimum requirement.

    Args:
        result: Integer count of unpacked extensions from getter
        expected: Dict with 'min_count' key containing the minimum expected count
        **options: Additional options (not used)

    Returns:
        float: 1.0 if count >= min_count, 0.0 otherwise
    """
    if not isinstance(result, int):
        return 0.0
    min_count = expected.get('min_count', 1)
    if result >= min_count:
        return 1.0
    else:
        return 0.0

def check_extension_dir_exists__42aa293803886f1362c8d3d5e33a2703(actual: dict, expected: dict, **options) -> float:
    """
    Check if VSCode extension directory exists.

    Args:
        actual (dict): Result from get_vscode_extension_dir getter with 'exists' field
        expected (dict): Rules dict with 'should_exist' field (rules are passed directly, not wrapped)
        **options: Additional options

    Returns:
        float: 1.0 if actual state matches expected rule, 0.0 otherwise
    """
    if not actual:
        return 0.0
    should_exist = expected.get('should_exist', True)
    actual_exists = actual.get('exists', False) if isinstance(actual, dict) else False
    if actual_exists == should_exist:
        return 1.0
    return 0.0

def is_expected_url_pattern_match__9aac585d873d78f321c8f29733249832(result, expected, **options) -> float:
    """
    Checks if the active tab URL matches multiple regex patterns.

    Args:
        result: The active tab info (string URL or dict with 'url' field)
        expected: Dictionary with 'expected' key containing list of regex patterns
        **options: Additional options

    Returns:
        float: 1.0 if all patterns match, 0.0 otherwise
    """
    if not result:
        logger.info('[PATTERN_MATCH] No result provided')
        return 0.0
    if isinstance(result, str):
        result_url = result
    elif isinstance(result, dict) and 'url' in result:
        result_url = result['url']
    else:
        logger.error(f'[PATTERN_MATCH] Invalid result format: {type(result)}')
        return 0.0
    logger.info(f'[PATTERN_MATCH] Result URL: {result_url}')
    patterns = expected.get('expected', [])
    logger.info(f'[PATTERN_MATCH] Expected patterns: {patterns}')
    for pattern in patterns:
        match = re.search(pattern, result_url, re.IGNORECASE)
        if not match:
            logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' not found in URL")
            return 0.0
        logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' matched")
    logger.info('[PATTERN_MATCH] All patterns matched successfully')
    return 1.0

def check_chrome_extension_manifest__61167ddb747c3a88a31446374f5a958f(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate Chrome extension manifest and files.

    Args:
        result: Dict from getter with 'manifest', 'files_exist', 'all_files_exist'
        expected: Dict with expected manifest structure rules:
            - manifest_checks: List of dicts with 'key' (list of nested keys) and 'value' (expected value)
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    manifest = result.get('manifest', {})
    all_files_exist = result.get('all_files_exist', False)
    file_score = 0.5 if all_files_exist else 0.0
    manifest_checks = expected.get('manifest_checks', [])
    if not manifest_checks:
        return file_score
    passed_checks = 0
    total_checks = len(manifest_checks)
    for check in manifest_checks:
        key_path = check.get('key', [])
        expected_value = check.get('value')
        current = manifest
        try:
            for key in key_path:
                current = current[key]
            if current == expected_value:
                passed_checks += 1
        except (KeyError, TypeError):
            pass
    manifest_score = 0.5 * (passed_checks / total_checks) if total_checks > 0 else 0.0
    return file_score + manifest_score

def check_single_tab_url__c557d7de5c301473caee703726a81950(result: List[Dict[str, str]], expected: Dict[str, Any], **options) -> float:
    """Check if only one tab is open with the expected URL.

    Args:
        result: List of open tabs from get_open_tabs_info
        expected: Rules dict containing expected_url and should_be_only_tab
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 with partial credit
    """
    expected_url = expected.get('expected_url', '')
    should_be_only_tab = expected.get('should_be_only_tab', False)
    if not isinstance(result, list):
        return 0.0
    result_urls = []
    for tab in result:
        if isinstance(tab, dict) and 'url' in tab:
            result_urls.append(tab['url'])

    def normalize_url(url):
        return url.rstrip('/')
    result_urls_normalized = [normalize_url(url) for url in result_urls]
    expected_url_normalized = normalize_url(expected_url)
    score = 0.0
    if expected_url_normalized in result_urls_normalized:
        score += 0.5
    if should_be_only_tab and len(result_urls) == 1:
        score += 0.5
    return min(score, 1.0)

def verify_tabs_with_state_change__58565672(open_tabs: List[Dict[str, str]], rule: Dict[str, Any]) -> float:
    """
    Verifies that:
    1. Expected tabs are present in final state (Scholar and Amazon)
    2. Apple tab has been closed/removed (present in initial, absent in final)
    3. Tab state actually changed from initial to final

    This function compares initial state (from postconfig) with final state to verify
    that the Apple tab was actually closed, not just never existed.

    Args:
        open_tabs: List of dicts containing final tab information (each with 'url' key)
        rule: Dict containing:
            - 'initial_tabs': List of dicts from postconfig showing initial tab state
            - 'expected_present': List of URLs that should be present in final state
            - 'expected_absent': List of URLs that should NOT be present in final state

    Returns:
        float: 1.0 if all conditions met, 0.0 otherwise
    """
    if not open_tabs:
        logger.error('No tabs found in final state')
        return 0.0
    expected_present = rule.get('expected_present', [])
    expected_absent = rule.get('expected_absent', [])
    initial_tabs = rule.get('initial_tabs', [])
    if not initial_tabs:
        logger.warning('No initial tabs data available - cannot verify state change')
        initial_tabs = []
    actual_urls = [tab['url'] for tab in open_tabs]
    initial_urls = [tab['url'] for tab in initial_tabs] if initial_tabs else []
    logger.info(f'Initial URLs: {initial_urls}')
    logger.info(f'Final URLs: {actual_urls}')
    logger.info(f'Expected present: {expected_present}')
    logger.info(f'Expected absent: {expected_absent}')

    def compare_urls(url1: str, url2: str) -> bool:
        """Compare two URLs, handling trailing slashes and common variations."""
        if url1 == url2:
            return True
        url1_norm = url1.rstrip('/')
        url2_norm = url2.rstrip('/')
        return url1_norm == url2_norm
    for expected_url in expected_present:
        if not any((compare_urls(expected_url, actual_url) for actual_url in actual_urls)):
            logger.error(f'Expected URL not found in final state: {expected_url}')
            return 0.0
    for absent_url in expected_absent:
        if any((compare_urls(absent_url, actual_url) for actual_url in actual_urls)):
            logger.error(f'URL should have been closed but is still present: {absent_url}')
            return 0.0
    if initial_urls:
        for absent_url in expected_absent:
            if not any((compare_urls(absent_url, initial_url) for initial_url in initial_urls)):
                logger.error(f'URL marked as absent was not in initial state - cannot verify closure: {absent_url}')
                return 0.0
            else:
                logger.info(f'Verified {absent_url} was in initial state and removed in final state')
    expected_count = len(expected_present)
    actual_count = len(actual_urls)
    if actual_count != expected_count:
        logger.error(f'Tab count mismatch. Expected: {expected_count}, Actual: {actual_count}')
        return 0.0
    if initial_urls and len(initial_urls) != 2:
        logger.warning(f'Initial tab count unexpected. Expected: 2, Actual: {len(initial_urls)}')
    logger.info('All tab verification checks passed - verified state change from initial to final')
    return 1.0

def check_table_count_only__e164fbd9(result, expected, **options):
    """
    Check if a new digraph table was added at the beginning.

    Verifies:
    1. Table count increased from 13 to 14
    2. First table is at the beginning of the document
    3. First table has exactly 1 row
    4. All cells in the first row contain digraphs (2-letter combinations)
    5. Has at least some digraphs (table is not empty)
    6. The digraphs in the table match the actual digraphs mentioned in the document
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.error('✗ Invalid result or expected format')
        return 0.0
    result_count = result.get('table_count', 0)
    expected_count = expected.get('table_count', 0)
    first_table_info = result.get('first_table_info')
    document_digraphs = result.get('document_digraphs', [])
    score = 0.0
    max_score = 6.0
    if result_count == expected_count:
        logger.info(f'✓ Table count matches: {result_count}')
        score += 1.0
    else:
        logger.info(f'✗ Table count mismatch: got {result_count}, expected {expected_count}')
        return 0.0
    if not first_table_info:
        logger.error('✗ No first table information available')
        return 0.0
    if first_table_info.get('position') == 'first':
        logger.info('✓ Table is at the beginning of the document')
        score += 1.0
    else:
        logger.info(f"✗ Table is not at the beginning: {first_table_info.get('position')}")
    if first_table_info.get('is_single_row'):
        logger.info(f'✓ Table has exactly 1 row')
        score += 1.0
    else:
        logger.info(f"✗ Table has {first_table_info.get('row_count')} rows, expected 1")
    if first_table_info.get('are_digraphs'):
        logger.info('✓ All cells contain 2-letter digraphs')
        score += 1.0
    else:
        cells = first_table_info.get('cells', [])
        logger.info(f'✗ Not all cells are digraphs. Cells: {cells}')
    cells = first_table_info.get('cells', [])
    non_empty_cells = [c for c in cells if c]
    if len(non_empty_cells) >= 3:
        logger.info(f'✓ Table contains {len(non_empty_cells)} digraphs')
        score += 1.0
    else:
        logger.info(f'✗ Table has too few digraphs: {non_empty_cells}')
    if document_digraphs:
        table_digraphs = set((cell.lower() for cell in non_empty_cells if cell))
        doc_digraphs_set = set((dg.lower() for dg in document_digraphs))
        matching_digraphs = table_digraphs.intersection(doc_digraphs_set)
        if table_digraphs:
            coverage = len(matching_digraphs) / len(table_digraphs)
        else:
            coverage = 0.0
        if doc_digraphs_set:
            representation = len(matching_digraphs) / len(doc_digraphs_set)
        else:
            representation = 0.0
        logger.info(f'Document digraphs found: {len(doc_digraphs_set)}')
        logger.info(f'Table digraphs: {sorted(table_digraphs)[:10]}')
        logger.info(f'Matching digraphs: {sorted(matching_digraphs)[:10]} ({len(matching_digraphs)} total)')
        logger.info(f'Coverage: {coverage:.2%} of table digraphs are from document')
        logger.info(f'Representation: {representation:.2%} of document digraphs are in table')
        if coverage >= 0.8:
            if representation >= 0.7:
                logger.info('✓ Excellent match: Table contains nearly all digraphs from document')
                score += 1.0
            elif representation >= 0.5:
                logger.info('✓ Very good match: Table contains most digraphs from document')
                score += 0.9
            elif representation >= 0.3:
                logger.info('✓ Good match: Table contains many digraphs from document')
                score += 0.7
            else:
                logger.info('✓ Acceptable: Table digraphs are real but incomplete')
                score += 0.5
        elif coverage >= 0.6:
            logger.info(f'⚠ Moderate match: {coverage:.0%} of table digraphs are from document')
            score += 0.3
        else:
            logger.info(f'✗ Poor match: Only {coverage:.0%} of table digraphs are from document')
            logger.info(f"   Many table digraphs don't appear in the source document")
    else:
        logger.info('⚠ No document digraphs found - cannot verify content match')
        score += 0.1
    final_score = score / max_score
    logger.info(f'Final score: {final_score:.2f} ({score:.1f}/{max_score} points)')
    return final_score

def check_extension_name__de268faa9ad661d5adcd42e1b7f775e2(result, expected, **options):
    """Check if expected extension name is in the list of installed extensions.

    Args:
        result: List of extension names from getter
        expected: Dict with 'extension_name' key containing the expected name
        **options: Additional options (not used)

    Returns:
        float: 1.0 if extension name found, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    expected_name = expected.get('extension_name', '')
    if not expected_name:
        return 0.0
    for name in result:
        if name.lower() == expected_name.lower():
            return 1.0
    return 0.0

def check_extension_folder_exists__ae6416e4(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if extension folder exists AND extension is loaded into Chrome.

    The task requires TWO steps:
    1. Extract the extension folder (verify folder exists)
    2. Load the unpacked extension into Chrome (verify in Chrome Preferences)

    Args:
        result: Dict with:
            - 'folder_exists': boolean - Whether the folder exists
            - 'extension_installed': boolean - Whether extension is loaded in Chrome
            - 'folder_path': str - Path to the folder
            - 'installed_path': str - Path in Chrome Preferences (if installed)
        expected: Rules dict (not used, both checks required)
        **options: Additional options

    Returns:
        1.0 if BOTH folder exists AND extension is loaded in Chrome, 0.0 otherwise
    """
    folder_exists = result.get('folder_exists', False)
    extension_installed = result.get('extension_installed', False)
    folder_path = result.get('folder_path', '')
    installed_path = result.get('installed_path', '')
    logger.info(f'Folder exists: {folder_exists} at {folder_path}')
    logger.info(f'Extension installed in Chrome: {extension_installed}')
    if installed_path:
        logger.info(f'Extension Chrome path: {installed_path}')
    if folder_exists and extension_installed:
        logger.info('SUCCESS: Extension folder extracted AND loaded into Chrome')
        return 1.0
    elif folder_exists and (not extension_installed):
        logger.warning('PARTIAL: Extension folder extracted but NOT loaded into Chrome')
        return 0.0
    elif not folder_exists and extension_installed:
        logger.warning("UNEXPECTED: Extension loaded in Chrome but folder doesn't exist")
        return 0.0
    else:
        logger.error('FAILURE: Extension folder not extracted and not loaded in Chrome')
        return 0.0

def check_pdf_file_and_tab__a606126a(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if PDF file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with pdf_path, expected_urls, and pdf_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f'PDF file exists: +0.5')
    else:
        logger.warning(f'PDF file does not exist')
    chrome_tabs = result.get('chrome_tabs', [])
    pdf_url_pattern = expected.get('pdf_url_pattern', '')
    pdf_opened = False
    for tab_url in chrome_tabs:
        if tab_url == pdf_url_pattern:
            pdf_opened = True
            break
    if pdf_opened:
        score += 0.5
        logger.info(f'PDF file opened in Chrome: +0.5')
    else:
        logger.warning(f'PDF file not found in Chrome tabs. Tabs: {chrome_tabs}')
    logger.info(f'Final score: {score}')
    return score

def check_chrome_setting__0717aaed(result, expected, **options):
    """
    Check if the Chrome Password saving setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_gdrive_files__87199839(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_html_export_and_chrome__d6487d33(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if HTML file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with html_path and html_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f'HTML file exists: +0.5')
    else:
        logger.warning(f'HTML file does not exist')
    chrome_tabs = result.get('chrome_tabs', [])
    html_url_pattern = expected.get('html_url_pattern', '')
    html_opened = False
    for tab_url in chrome_tabs:
        if html_url_pattern in tab_url:
            html_opened = True
            break
    if html_opened:
        score += 0.5
        logger.info(f'HTML file opened in Chrome: +0.5')
    else:
        logger.warning(f'HTML file not found in Chrome tabs. Tabs: {chrome_tabs}')
    logger.info(f'Final score: {score}')
    return score

def check_gdrive_pdf__7e0bdf48(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def is_expected_bookmarks__7a5a7856_aug_6_task_verify_0(bookmarks: Dict[str, Any], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are present in Chrome bookmark bar.

    Uses subset checking to verify that all expected URLs are present in the bookmark bar,
    allowing for additional bookmarks to exist without failing the evaluation.

    This fixes the semantic mismatch where the instruction says "Bookmark the article"
    (add requirement) but the metric should not require exact match (which would fail
    if user had pre-existing bookmarks or bookmarked additional items).

    Args:
        bookmarks: Dict containing bookmark data from Chrome
        rule: Dict with 'type' and expected values

    Returns:
        float: 1.0 if all expected bookmarks are present, 0.0 otherwise
    """
    if not bookmarks:
        logger.info('No bookmarks data provided')
        return 0.0
    if rule['type'] == 'bookmark_bar_websites_urls':
        bookmark_bar_websites_urls = [bookmark['url'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'url']
        expected_urls = rule['urls']
        logger.info(f'Expected URLs: {expected_urls}')
        logger.info(f'Actual bookmark bar URLs: {bookmark_bar_websites_urls}')
        if set(expected_urls).issubset(set(bookmark_bar_websites_urls)):
            logger.info('All expected bookmarks are present in bookmark bar')
            return 1.0
        else:
            logger.info('Not all expected bookmarks are present in bookmark bar')
            return 0.0
    else:
        logger.error(f"Unknown rule type: {rule['type']}")
        return 0.0

def check_gdrive_file__3ce93f6ce10147dcc2489b34874a9f3b(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if the Google Drive file meets expected criteria.

    Args:
        result: Dict from getter with keys: exists, file_count, file_name
        expected: Dict with validation rules:
            - exists: bool - expected existence
            - min_count: int (optional) - minimum number of files expected
            - file_name: str (optional) - expected file name (exact match)

    Returns:
        float: 1.0 if all criteria met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_exists = expected.get('exists', True)
    if result.get('exists', False) != expected_exists:
        return 0.0
    if not expected_exists:
        return 1.0
    min_count = expected.get('min_count')
    if min_count is not None:
        if result.get('file_count', 0) < min_count:
            return 0.0
    expected_name = expected.get('file_name')
    if expected_name is not None:
        if result.get('file_name') != expected_name:
            return 0.0
    return 1.0

def check_gdrive_eml_count__26d2516b888edf7c1b328cca7acaf9b7(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if the correct number of .eml files exist in Google Drive folder.

    Args:
        result: List of .eml filenames from getter
        expected: Expected validation rules with:
            - expected_count: Expected number of .eml files
            - required_extension: File extension to verify (default: '.eml')

    Returns:
        1.0 if count matches expected, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_count = expected.get('expected_count', 2)
    required_extension = expected.get('required_extension', '.eml')
    valid_files = [f for f in result if f.endswith(required_extension)]
    if len(valid_files) == expected_count:
        logger.info(f'Found {len(valid_files)} {required_extension} files as expected')
        return 1.0
    else:
        logger.info(f'Expected {expected_count} {required_extension} files, found {len(valid_files)}')
        return 0.0

def check_webext_manifest__951cee96(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_gdrive_files__b0d15a73736689e26424d81f759f1528(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected files exist in Google Drive folder.

    Args:
        result: List of filenames found in Google Drive folder
        expected: Dict with 'expected_files' key containing list of required filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.warning('No expected files specified')
        return 0.0
    if not result:
        logger.warning('No files found in Google Drive')
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
            logger.info(f'Found expected file: {expected_file}')
        else:
            logger.warning(f'Missing expected file: {expected_file}')
    score = found_count / len(expected_files)
    logger.info(f'Found {found_count}/{len(expected_files)} expected files. Score: {score}')
    return score

def is_expected_bookmarks__7a5a7856_f1b6_42a4_ade9_1ca81ca0f263_task_verify_0(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are present in Chrome.
    Uses subset checking to verify that the expected URLs are present,
    without requiring exact match (allows pre-existing bookmarks).

    This is different from the original is_expected_bookmarks which uses
    exact set equality and would fail if there are any other bookmarks present.

    Args:
        bookmarks: Bookmarks data structure from Chrome
        rule: Rule dictionary containing type and expected values

    Returns:
        1.0 if expected bookmarks are present, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'bookmark_bar_websites_urls':
        bookmark_bar_websites_urls = [bookmark['url'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'url']
        expected_urls = set(rule['urls'])
        actual_urls = set(bookmark_bar_websites_urls)
        logger.info(f'Expected URLs: {expected_urls}')
        logger.info(f'Actual URLs in bookmark bar: {actual_urls}')
        if expected_urls.issubset(actual_urls):
            logger.info('All expected URLs are present in bookmark bar')
            return 1.0
        else:
            missing_urls = expected_urls - actual_urls
            logger.info(f'Missing URLs from bookmark bar: {missing_urls}')
            return 0.0
    else:
        logger.error(f"Unsupported rule type: {rule['type']}")
        return 0.0

def check_html_file_and_tab__3cb1a587(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if HTML file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with html_path, expected_urls, and html_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f'HTML file exists: +0.5')
    else:
        logger.warning(f'HTML file does not exist')
    chrome_tabs = result.get('chrome_tabs', [])
    html_url_pattern = expected.get('html_url_pattern', '')
    html_opened = False
    for tab_url in chrome_tabs:
        if html_url_pattern and html_url_pattern in tab_url:
            html_opened = True
            break
    if html_opened:
        score += 0.5
        logger.info(f'HTML file opened in Chrome: +0.5')
    else:
        logger.warning(f'HTML file not found in Chrome tabs. Expected pattern: {html_url_pattern}, Tabs: {chrome_tabs}')
    logger.info(f'Final score: {score}')
    return score

def check_website_filled__e7fe8e90(result, expected, **options):
    """Check if cell(s) are filled with valid data

    Args:
        result: Data from getter function
        expected: Expected rules dict
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    if isinstance(result, list):
        min_length = expected.get('min_length', 1)
        cells = expected.get('cells', [])
        if len(result) != len(cells):
            return 0.0
        score = 0.0
        for value in result:
            if value and len(value) >= min_length:
                score += 1.0 / len(result)
        return score
    elif isinstance(result, dict):
        name = result.get('name', '')
        address = result.get('address', '')
        min_name_length = expected.get('min_name_length', 1)
        min_address_length = expected.get('min_address_length', 1)
        score = 0.0
        if name and len(name) >= min_name_length:
            score += 0.5
        if address and len(address) >= min_address_length:
            score += 0.5
        return score
    elif isinstance(result, int):
        min_count = expected.get('min_count', 1)
        if result >= min_count:
            return 1.0
        else:
            return result / min_count if min_count > 0 else 0.0
    elif 'must_contain' in expected:
        if expected['must_contain'].lower() in result.lower():
            return 1.0
        else:
            return 0.0
    else:
        min_length = expected.get('min_length', 1)
        if result and len(result) >= min_length:
            return 1.0
        else:
            return 0.0

def check_chrome_extension_manifest__70453fbf044a41274b6b7dc7e909fb11(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate Chrome extension manifest and files.

    Args:
        result: Dict from getter with 'manifest', 'files_exist', 'all_files_exist'
        expected: Dict with expected manifest structure rules:
            - manifest_checks: List of dicts with 'key' (list of nested keys) and 'value' (expected value)
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    manifest = result.get('manifest', {})
    all_files_exist = result.get('all_files_exist', False)
    file_score = 0.5 if all_files_exist else 0.0
    manifest_checks = expected.get('manifest_checks', [])
    if not manifest_checks:
        return file_score
    passed_checks = 0
    total_checks = len(manifest_checks)
    for check in manifest_checks:
        key_path = check.get('key', [])
        expected_value = check.get('value')
        current = manifest
        try:
            for key in key_path:
                current = current[key]
            if current == expected_value:
                passed_checks += 1
        except (KeyError, TypeError):
            if expected_value == '':
                passed_checks += 1
            pass
    manifest_score = 0.5 * (passed_checks / total_checks) if total_checks > 0 else 0.0
    return file_score + manifest_score

def check_webext_manifest__4b1e45b68c5d85573f69177601102dcb(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if web extension manifest.json matches expected structure.

    Args:
        result: Manifest JSON dict from getter (or None if file not found)
        expected: Expected manifest structure with rules to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.warning('Manifest file not found')
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_options_page = expected.get('has_options_page', False)
    expected_no_description = expected.get('no_description', False)
    expected_no_page_action = expected.get('no_page_action', False)
    score = 0.0
    total_checks = 0
    if expected_name is not None:
        total_checks += 1
        if result.get('name') == expected_name:
            score += 1
            logger.info(f'Name matches: {expected_name}')
        else:
            logger.warning(f"Name mismatch: expected '{expected_name}', got '{result.get('name')}'")
    if expected_version is not None:
        total_checks += 1
        if result.get('version') == expected_version:
            score += 1
            logger.info(f'Version matches: {expected_version}')
        else:
            logger.warning(f"Version mismatch: expected '{expected_version}', got '{result.get('version')}'")
    if expected_has_background:
        total_checks += 1
        if 'background' in result and 'scripts' in result['background'] and (len(result['background']['scripts']) > 0):
            score += 1
            logger.info(f"Background scripts found: {result['background']['scripts']}")
        else:
            logger.warning('Background scripts not found')
    if expected_has_browser_action:
        total_checks += 1
        if 'browser_action' in result:
            score += 1
            logger.info('Browser action found')
        else:
            logger.warning('Browser action not found')
    if expected_has_content_scripts:
        total_checks += 1
        if 'content_scripts' in result and len(result['content_scripts']) > 0:
            score += 1
            logger.info(f"Content scripts found: {len(result['content_scripts'])} entries")
        else:
            logger.warning('Content scripts not found')
    if expected_has_options_page:
        total_checks += 1
        if 'options_page' in result or 'options_ui' in result:
            score += 1
            logger.info('Options page found')
        else:
            logger.warning('Options page not found')
    if expected_no_description:
        total_checks += 1
        description = result.get('description', '')
        if not description or description.strip() == '':
            score += 1
            logger.info('Description is blank as expected')
        else:
            logger.warning(f"Description should be blank but found: '{description}'")
    if expected_no_page_action:
        total_checks += 1
        if 'page_action' not in result:
            score += 1
            logger.info('page_action not present as expected')
        else:
            logger.warning('page_action should not be included but was found')
    if total_checks == 0:
        logger.warning('No checks to perform')
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks} checks passed)')
    return final_score

def check_eslint_extension__d09f7694(actual: str, rules: Dict, **options):
    """
    Check if ESLint extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_table_top_left__2d408821(result, expected, **options):
    """Check if table is positioned in top-left corner.

    Args:
        result: Table position dict from getter
        expected: Dict with 'max_top' and 'max_left'
        **options: Additional options

    Returns:
        float: 1.0 if both conditions met, 0.0 otherwise
    """
    if result is None:
        return 0.0
    max_top = expected.get('max_top', 0)
    max_left = expected.get('max_left', 0)
    if result['top'] < max_top and result['left'] < max_left:
        return 1.0
    else:
        return 0.0

def check_extension_manifest__3d9fc0c33aef9d5f768043f561973d63(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate browser extension manifest.json structure.

    Args:
        result: Manifest data from getter
        expected: Rules dict containing expected values for name, version, background, browser_action
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct field)
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dictionary')
        return 0.0
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    check_background = expected.get('check_background', False)
    check_browser_action = expected.get('check_browser_action', False)
    score = 0.0
    total_checks = 0
    total_checks += 1
    if result.get('name') == expected_name:
        score += 0.25
        logger.info(f'Name matches: {expected_name}')
    else:
        logger.info(f"Name mismatch: expected={expected_name}, got={result.get('name')}")
    total_checks += 1
    if result.get('version') == expected_version:
        score += 0.25
        logger.info(f'Version matches: {expected_version}')
    else:
        logger.info(f"Version mismatch: expected={expected_version}, got={result.get('version')}")
    if check_background:
        total_checks += 1
        background = result.get('background', {})
        if isinstance(background, dict) and 'scripts' in background:
            scripts = background.get('scripts', [])
            if isinstance(scripts, list) and len(scripts) > 0:
                score += 0.25
                logger.info(f'Background scripts found: {scripts}')
            else:
                logger.info('Background scripts field exists but is empty')
        else:
            logger.info('Background section missing or invalid')
    if check_browser_action:
        total_checks += 1
        browser_action = result.get('browser_action', {})
        if isinstance(browser_action, dict) and len(browser_action) > 0:
            has_fields = any((k in browser_action for k in ['default_popup', 'default_icon', 'default_title']))
            if has_fields:
                score += 0.25
                logger.info(f'Browser action found with fields: {list(browser_action.keys())}')
            else:
                logger.info('Browser action exists but missing expected fields')
        else:
            logger.info('Browser action section missing or invalid')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_bookmark_folder_contains_url__3c0977b51a553edc4a7433583af5fcf6(bookmarks: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Checks if a specific folder exists in the bookmark bar and contains a specific URL.

    Args:
        bookmarks: Bookmarks data structure from get_bookmarks getter
        expected: Expected configuration with:
            - folder_name (str): Name of the folder to check
            - url (str): URL that should be in the folder
        **options: Additional options (unused)

    Returns:
        float: 0.5 if folder exists, 1.0 if folder exists and contains URL, 0.0 otherwise
    """
    if not bookmarks:
        logger.info('No bookmarks data available')
        return 0.0
    folder_name = expected.get('folder_name')
    expected_url = expected.get('url')
    if not folder_name or not expected_url:
        logger.error(f'Missing required parameters: folder_name={folder_name}, url={expected_url}')
        return 0.0
    bookmark_bar = bookmarks.get('bookmark_bar', {})
    children = bookmark_bar.get('children', [])
    target_folder = None
    for bookmark in children:
        if bookmark.get('type') == 'folder' and bookmark.get('name') == folder_name:
            target_folder = bookmark
            break
    if not target_folder:
        logger.info(f"Folder '{folder_name}' not found in bookmark bar")
        return 0.0
    logger.info(f"Folder '{folder_name}' found in bookmark bar")
    folder_children = target_folder.get('children', [])
    folder_urls = [child.get('url') for child in folder_children if child.get('type') == 'url']
    logger.info(f"URLs in folder '{folder_name}': {folder_urls}")
    if expected_url in folder_urls:
        logger.info(f"URL '{expected_url}' found in folder '{folder_name}'")
        return 1.0
    else:
        logger.info(f"URL '{expected_url}' not found in folder '{folder_name}'")
        return 0.5

def is_expected_bookmarks__7a5a7856_aug_6(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome using subset logic.

    For 'Add' operations, this verifies that all expected URLs are present
    in the bookmark bar, allowing for additional pre-existing bookmarks.

    Args:
        bookmarks: Chrome bookmarks data structure
        rule: Rule specification with 'type' and 'urls' keys

    Returns:
        float: 1.0 if all expected URLs are present (subset), 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'bookmark_bar_websites_urls':
        bookmark_bar_websites_urls = [bookmark['url'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'url']
        expected_urls = set(rule['urls'])
        actual_urls = set(bookmark_bar_websites_urls)
        if expected_urls.issubset(actual_urls):
            return 1.0
        else:
            return 0.0
    return 0.0

def check_partial_extensions__2f01adaf(result, expected, **options):
    """
    Check installed extensions with partial credit scoring.
    Each correctly installed extension contributes to the score.

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' list of extension names
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on proportion of expected extensions installed
    """
    if not result:
        result = []
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        return 0.0
    set_expected = set(expected_extensions)
    set_installed = set(result)
    matched = set_expected.intersection(set_installed)
    score = len(matched) / len(set_expected)
    logger.info(f'Expected extensions: {expected_extensions}')
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Matched: {matched}, Score: {score}')
    return score

def check_html_file_exists__48e4da460d6f3e132e4d3cc48ac9ca24(result, expected, **options):
    """
    Check if HTML file exists.

    Args:
        result: Dict from getter with 'exists' key
        expected: Dict with expected values (from rules)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Invalid result type: {type(result)}')
        return 0.0
    exists = result.get('exists', False)
    if exists:
        logger.info(f"✅ HTML file exists at {result.get('path', 'unknown')}")
        return 1.0
    else:
        logger.warning(f"❌ HTML file does not exist at {result.get('path', 'unknown')}")
        return 0.0

def check_extension_name__a4b0d616259699774cd5ee2679c8a96c(result, expected, **options):
    """
    Check if the extension name matches the expected value.

    Args:
        result: String from getter containing the extension name
        expected: Dict with 'extension_name' key
        **options: Additional options (unused)

    Returns:
        float: 1.0 if names match, 0.0 otherwise
    """
    expected_name = expected.get('extension_name', '')
    if result == expected_name:
        return 1.0
    return 0.0

def check_gdrive_pdf_info__bb91e7693f2a30704f1d1cc79be73950(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if PDF exists in Google Drive with expected properties.

    Args:
        result: Dictionary from getter with 'exists', 'page_count', 'file_size' keys
        expected: Expected rules dict with 'page_count', 'min_size' keys
        **options: Additional options

    Returns:
        float: Score from 0.0 to 1.0 based on verification results
    """
    if result is None or not result.get('exists', False):
        logger.info('PDF file not found in Google Drive')
        return 0.0
    score = 0.0
    expected_pages = expected.get('page_count', 1)
    actual_pages = result.get('page_count', 0)
    if actual_pages == expected_pages:
        score += 0.5
        logger.info(f'Page count matches: {actual_pages} pages')
    else:
        logger.info(f'Page count mismatch: expected {expected_pages}, got {actual_pages}')
    min_size = expected.get('min_size', 5120)
    actual_size = result.get('file_size', 0)
    if actual_size >= min_size:
        score += 0.5
        logger.info(f'File size acceptable: {actual_size} bytes >= {min_size} bytes')
    else:
        logger.info(f'File size too small: {actual_size} bytes < {min_size} bytes')
    return score

def check_chrome_bookmark_exists__b5b6115e(result: list, expected: dict, **options):
    """Check if bookmark exists with URL and title.

    Args:
        result: List of bookmarks from getter
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_url = expected.get('url', 'arxiv.org')
    expected_title = expected.get('title', 'ArXiv Papers')
    for bookmark in result:
        if isinstance(bookmark, dict):
            url = bookmark.get('url', '')
            title = bookmark.get('name', '')
            if expected_url in url and expected_title in title:
                logger.info(f'Found matching bookmark: {title} - {url}')
                return 1.0
    logger.info(f"Bookmark with URL '{expected_url}' and title '{expected_title}' not found")
    return 0.0

def check_gdrive_pdf_in_folder__661594c9(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_author_table_structure__f918ebce432b3a8956b1f7dd26a64d11(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if the table has correct structure, headers, row count, content validation, and actual author data verification.

    Args:
        result: Table data from getter with 'headers', 'rows', 'row_count', 'data'
        expected: Expected data with 'required_headers', 'min_columns', 'expected_row_count', 'sample_authors'
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    score = 0.0
    headers = result.get('headers', [])
    rows = result.get('rows', [])
    row_count = result.get('row_count', 0)
    data = result.get('data', [])
    required_headers = expected.get('required_headers', [])
    min_columns = expected.get('min_columns', 3)
    expected_row_count = expected.get('expected_row_count', 4)
    sample_authors = expected.get('sample_authors', [])
    if row_count == expected_row_count:
        score += 0.15
    elif row_count > 0:
        score += 0.075
    if headers:
        headers_normalized = [h.lower().strip() for h in headers if h]
        required_normalized = [h.lower().strip() for h in required_headers]
        headers_match = 0
        for req_header in required_normalized:
            if req_header in headers_normalized:
                headers_match += 1
        if headers_match == len(required_normalized):
            score += 0.2
        else:
            score += 0.2 * (headers_match / len(required_normalized))
    if rows:
        valid_rows = sum((1 for row in rows if len(row) >= min_columns))
        if valid_rows == len(rows):
            score += 0.1
        else:
            score += 0.1 * (valid_rows / len(rows))
    email_score = _validate_email_format(data, headers)
    score += email_score * 0.15
    affiliation_score = _validate_affiliations(data, headers)
    score += affiliation_score * 0.15
    name_score = _validate_name_format(data, headers)
    score += name_score * 0.1
    content_score = _verify_sample_authors(data, headers, sample_authors)
    score += content_score * 0.15
    return min(score, 1.0)

def check_java_extension__c2202196(actual: str, rules: Dict, **options):
    """
    Check if Java Extension Pack is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_googledrive_files_exist__ba61b137046435f47239f8911466a875(result: Optional[List[str]], expected, **options):
    """Check if expected files exist in Google Drive folder.

    Args:
        result: List of filenames from getter, or None if folder doesn't exist
        expected: Dict with 'required_files' list
        **options: Additional options

    Returns:
        float: 1.0 if all required files exist, 0.0 otherwise
    """
    if result is None:
        logger.info('Result is None - folder not found')
        return 0.0
    required_files = expected.get('required_files', [])
    if not required_files:
        logger.warning('No required_files specified in expected')
        return 0.0
    result_set = set(result)
    missing_files = [f for f in required_files if f not in result_set]
    if missing_files:
        logger.info(f'Missing files: {missing_files}')
        logger.info(f'Found files: {result}')
        return 0.0
    logger.info(f'All required files found: {required_files}')
    return 1.0

def check_extension_not_installed__cd5e025975a2ca9b4b85d9a92890bab8(actual: str, expected: dict, **options) -> float:
    """
    Check if the VSCode extension is NOT installed (for uninstall verification).

    The extension is pre-installed in the main config (before task execution) to ensure it
    exists when the agent starts. The agent must uninstall it. This prevents false positives
    where an agent that does nothing would pass.

    Args:
        actual (str): Output from 'code --list-extensions' command
        expected (dict): Expected rules with 'extension_id' field
        **options: Additional options

    Returns:
        float: 1.0 if extension is NOT installed (uninstall succeeded), 0.0 if it is still installed (uninstall failed)
    """
    if not actual:
        return 1.0
    extension_id = expected.get('extension_id', '')
    if not extension_id:
        return 0.0
    installed_extensions = [line.strip() for line in actual.strip().split('\n') if line.strip()]
    for ext in installed_extensions:
        if ext.lower() == extension_id.lower():
            return 0.0
    return 1.0

def check_last_table_structure__e0b81f05(result, expected, **options):
    """Check last table structure with content verification."""
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    score = 0.0
    min_rows = expected.get('min_rows', 0)
    if result.get('row_count', 0) >= min_rows:
        score += 0.25
    else:
        return 0.0
    expected_cols = expected.get('columns', 0)
    if result.get('col_count', 0) == expected_cols:
        score += 0.25
    else:
        return 0.0
    headers = result.get('headers', [])
    header_text = ' '.join(headers).lower()
    if 'type' in header_text and 'description' in header_text and ('example' in header_text):
        score += 0.25
    else:
        return score
    rows = result.get('rows', [])
    required_types = ['graph', 'digraph', 'trigraph', 'quadgraph']
    all_row_text = ' '.join([' '.join(row).lower() for row in rows])
    types_found = sum((1 for req_type in required_types if req_type in all_row_text))
    if types_found == len(required_types):
        score += 0.25
    else:
        score += 0.25 * (types_found / len(required_types))
    return score

def is_extension_installed__f0c7e5c8(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_gdrive_pdf__a70aba9d(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_url_deleted__b835bc31(history_data: List, rule: Dict[str, Any]) -> float:
    """
    Check if specific URLs are deleted from Chrome history.

    Args:
        history_data: List of tuples (url, title, last_visit_time)
        rule: Dict with 'type' and 'urls' keys
            - type: must be 'urls'
            - urls: list of URLs that should be deleted

    Returns:
        1.0 if all specified URLs are deleted, 0.0 otherwise
    """
    if rule['type'] == 'urls':
        history_urls = [history[0] for history in history_data]
        for target_url in rule['urls']:
            if target_url in history_urls:
                logger.info(f'URL still exists in history: {target_url}')
                return 0.0
        logger.info(f'All target URLs successfully deleted')
        return 1.0
    else:
        raise TypeError(f"{rule['type']} not supported yet!")

def check_gdrive_pdf__bf7218b9(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_url_replaced_with_doi__e1878192(result, expected, **options):
    """Check if URL was replaced with DOI in reference.

    Args:
        result: Path to the result DOCX file
        expected: Dict with 'author', 'year', 'should_not_contain', 'doi' keys

    Returns:
        float: Score for URL removal and DOI addition
    """
    if not result or not isinstance(result, str):
        return 0.0
    try:
        doc = Document(result)
        author = expected.get('author', '')
        year = expected.get('year', '')
        should_not_contain = expected.get('should_not_contain', '')
        doi = expected.get('doi', '')
        in_references = False
        target_para = None
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if 'references' in para_text.lower() and len(para_text) < 50:
                in_references = True
                continue
            if in_references and author in para_text and (year in para_text):
                target_para = para_text
                break
        if not target_para:
            return 0.0
        score = 0.0
        if should_not_contain not in target_para:
            score += 0.5
        if doi in target_para:
            score += 0.5
        return score
    except Exception as e:
        print(f'Error in check_url_replaced_with_doi__e1878192: {e}')
        return 0.0

def check_search_engine_changed__8a4fe4a4(result_state: str, expected: Dict[str, str], **options) -> float:
    """
    Check if the default search engine was changed from old_pattern to new_pattern.

    Args:
        result_state: The template_url string from Chrome preferences
        expected: Dictionary containing 'old_pattern' and 'new_pattern' (e.g., {'old_pattern': 'yahoo', 'new_pattern': 'google'})
        **options: Additional options (not used)

    Returns:
        float: 1.0 if the search engine was changed correctly (old_pattern removed AND new_pattern present), 0.0 otherwise
    """
    if not result_state:
        logger.warning('Result state is empty')
        return 0.0
    old_pattern = expected.get('old_pattern', '').lower()
    new_pattern = expected.get('new_pattern', '').lower()
    result_lower = result_state.lower()
    logger.info(f'Checking search engine change:')
    logger.info(f'  Template URL: {result_state}')
    logger.info(f'  Old pattern (should be removed): {old_pattern}')
    logger.info(f'  New pattern (should be present): {new_pattern}')
    if old_pattern and old_pattern in result_lower:
        logger.info(f"FAIL: Old pattern '{old_pattern}' still present in template URL")
        return 0.0
    if new_pattern and new_pattern not in result_lower:
        logger.info(f"FAIL: New pattern '{new_pattern}' not found in template URL")
        return 0.0
    logger.info('PASS: Search engine successfully changed from old pattern to new pattern')
    return 1.0

def check_table_height__60d42be3(result, expected, **options):
    """Check if table height meets minimum requirement.

    Args:
        result: Table position dict from getter
        expected: Dict with 'min_height' threshold
        **options: Additional options

    Returns:
        float: 1.0 if table height >= min_height, 0.0 otherwise
    """
    if result is None:
        return 0.0
    min_height = expected.get('min_height', 0)
    if result['height'] >= min_height:
        return 1.0
    else:
        return 0.0

def check_gdrive_files__5876322dac1e30402171f6bcd2edb019(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected files exist in Google Drive folder.

    Args:
        result: List of filenames found in Google Drive folder
        expected: Dict with 'expected_files' key containing list of required filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.warning('No expected files specified')
        return 0.0
    if not result:
        logger.warning('No files found in Google Drive')
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
            logger.info(f'Found expected file: {expected_file}')
        else:
            logger.warning(f'Missing expected file: {expected_file}')
    score = found_count / len(expected_files)
    logger.info(f'Found {found_count}/{len(expected_files)} expected files. Score: {score}')
    return score

def check_pdf_in_gdrive__a1937ab0(result: str, expected: Dict[str, Any], **options) -> float:
    """
    Check if PDF file exists in Google Drive and meets basic requirements.

    Args:
        result: Local path to downloaded PDF file from Google Drive
        expected: Dict with 'rules' key containing validation criteria
        **options: Additional options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    min_pages = expected.get('min_pages', 1)
    min_file_size = expected.get('min_file_size', 0)
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        file_size = os.path.getsize(result)
        if file_size < min_file_size:
            return 0.0
        reader = PdfReader(result)
        num_pages = len(reader.pages)
        if num_pages < min_pages:
            return 0.0
        return 1.0
    except Exception as e:
        return 0.0

def check_bookmark_folder_with_url__7a5a7856f1b642a4ade91ca81ca0f263000020251221151547(bookmarks: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a bookmark folder with a specific name exists in the bookmark bar
    and contains a specific URL.

    Args:
        bookmarks: Bookmark data structure from get_bookmarks
        expected: Dict with 'folder_name' and 'url' keys

    Returns:
        float: 1.0 if folder exists and contains URL, 0.0 otherwise
    """
    if not bookmarks:
        logger.info('No bookmarks data available')
        return 0.0
    folder_name = expected.get('folder_name')
    target_url = expected.get('url')
    if not folder_name or not target_url:
        logger.error('Missing folder_name or url in expected config')
        return 0.0
    logger.info(f"Looking for folder '{folder_name}' containing URL '{target_url}'")
    bookmark_bar = bookmarks.get('bookmark_bar', {})
    children = bookmark_bar.get('children', [])
    target_folder = None
    for bookmark in children:
        if bookmark.get('type') == 'folder' and bookmark.get('name') == folder_name:
            target_folder = bookmark
            logger.info(f"Found folder '{folder_name}'")
            break
    if not target_folder:
        logger.info(f"Folder '{folder_name}' not found in bookmark bar")
        return 0.0
    folder_children = target_folder.get('children', [])
    folder_urls = [child.get('url') for child in folder_children if child.get('type') == 'url']
    logger.info(f"URLs in folder '{folder_name}': {folder_urls}")
    if target_url in folder_urls:
        logger.info(f"URL '{target_url}' found in folder '{folder_name}'")
        return 1.0
    else:
        logger.info(f"URL '{target_url}' not found in folder '{folder_name}'")
        return 0.0

def check_gitlens_extension__136a9bb5(actual: str, rules: Dict, **options):
    """
    Check if GitLens extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_vim_tabstop_enabled__76faceaeb09651393a582783b81e5798(result, expected, **options):
    """Check if the vim tabstop configuration output matches expected.

    Args:
        result: Output from vim config check command
        expected: Expected configuration strings (dict with include/exclude lists)
        **options: Additional options

    Returns:
        float: 1.0 if all expected strings are present and no excluded strings, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_gdrive_pdf_in_folder__ada1e84e(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def is_extension_installed__fdadaf5d(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_chrome_ext_single__abfa4b043befc7603aff32afef71bd1b(result: List[str], expected: Dict, **options) -> float:
    """Check if expected Chrome extension is installed.

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' key containing list of required extension names
        **options: Additional options (not used)

    Returns:
        1.0 if all expected extensions are installed, 0.0 otherwise
    """
    if not result:
        logger.info('No extensions found in result')
        return 0.0
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        logger.warning('No expected extensions specified')
        return 0.0
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Expected extensions: {expected_extensions}')
    set_expected = set(expected_extensions)
    set_installed = set(result)
    if set_expected.issubset(set_installed):
        logger.info('All expected extensions are installed')
        return 1.0
    else:
        missing = set_expected - set_installed
        logger.info(f'Missing extensions: {missing}')
        return 0.0

def check_gdrive_files__12a331f2(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_table_top_position__5ff0bd58(result, expected, **options):
    """Check if table has moved upward (smaller top value).

    Args:
        result: Table position dict from getter
        expected: Dict with 'max_top' threshold
        **options: Additional options

    Returns:
        float: 1.0 if table top < max_top, 0.0 otherwise
    """
    if result is None:
        return 0.0
    max_top = expected.get('max_top', 0)
    if result['top'] < max_top:
        return 1.0
    else:
        return 0.0

def check_html_files_exist__98346f3b(result, expected, **options):
    """Check if expected HTML files exist with correct content.

    Args:
        result: Dict with:
            - count: Number of HTML files
            - files: List of dicts with filename and content info
        expected: Dict with:
            - expected_count: Expected number of HTML files
            - filename_patterns: List of expected filename patterns (optional)
            - content_markers: List of expected content markers (optional)

    Returns:
        Score between 0.0 and 1.0
    """
    expected_count = expected.get('expected_count', 2)
    filename_patterns = expected.get('filename_patterns', ['agent', 'human-data-quality'])
    content_markers = expected.get('content_markers', ['Lilian Weng', 'LLM Powered Autonomous Agents', 'Human Data Quality'])
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {type(result)}')
        return 0.0
    if result.get('count', 0) != expected_count:
        logger.error(f"Expected {expected_count} files, got {result.get('count', 0)}")
        return 0.0
    files = result.get('files', [])
    if len(files) != expected_count:
        logger.error(f'Expected {expected_count} files in list, got {len(files)}')
        return 0.0
    has_agent_post = False
    has_human_data_quality_post = False
    for file_info in files:
        filename = file_info.get('filename', '').lower()
        if ('agent' in filename or file_info.get('has_agent_keyword', False)) and (file_info.get('has_2023_06_23', False) or '2023' in filename):
            if file_info.get('has_lilian_weng', False):
                has_agent_post = True
        if (file_info.get('has_human_data_quality', False) or 'human' in filename or 'data' in filename or ('quality' in filename)) and (file_info.get('has_2024_02_05', False) or '2024' in filename):
            if file_info.get('has_lilian_weng', False):
                has_human_data_quality_post = True
    if has_agent_post and has_human_data_quality_post:
        return 1.0
    else:
        logger.error(f'Missing required posts. Agent post: {has_agent_post}, Human data quality post: {has_human_data_quality_post}')
        return 0.0

def is_extension_installed__ba154228(actual: str, expected: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        expected: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if expected['type'] == 'contain':
        if expected['expected'] in actual:
            return 1.0
        return 0.0
    elif expected['type'] == 'not_contain':
        if expected['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {expected['type']}")

def check_csv_table_biz__fb9ae858(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_gdrive_pdf__82f685a6(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_chrome_experiments_exact_match__d7481e87a8a6dde50669ce517c215edf(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the enabled Chrome experiments exactly match the expected list.

    Args:
        result: List of enabled experiment names
        expected: Dict with 'experiment_names' key containing the exact list of enabled experiments
        **options: Additional options

    Returns:
        1.0 if the experiments exactly match, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    expected_experiments = expected.get('experiment_names', [])
    if not isinstance(expected_experiments, list):
        logger.error(f'experiment_names is not a list: {type(expected_experiments)}')
        return 0.0
    result_sorted = sorted(result)
    expected_sorted = sorted(expected_experiments)
    if result_sorted == expected_sorted:
        logger.info(f'Experiments exactly match: {expected_sorted}')
        return 1.0
    else:
        logger.info(f"Experiments don't match. Expected: {expected_sorted}, Got: {result_sorted}")
        return 0.0

def check_table_count__4217a000(docx_file, expected, **options):
    """Check if the document has the expected number of tables."""
    if not docx_file:
        return 0.0
    try:
        doc = Document(docx_file)
        actual_count = len(doc.tables)
        return 1.0 if actual_count == expected else 0.0
    except Exception as e:
        logger.error(f'Error: {e}')
        return 0.0

def is_expected_bookmarks__2ad9387a_65d8_4e33_ad5b_7580065a27ca_task_verify_3(bookmarks: Dict[str, Any], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.

    This custom metric uses subset checking for bookmark_bar_folders_names type,
    which verifies that the specified folder(s) exist on the bookmark bar,
    without requiring an exact match (other folders can also exist).

    Args:
        bookmarks: Bookmarks data from Chrome
        rule: Rule dict with 'type' and type-specific fields (e.g., 'names' for folders)

    Returns:
        float: 1.0 if the expected bookmarks/folders exist, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'bookmark_bar_folders_names':
        try:
            bookmark_bar_folders_names = [bookmark['name'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder']
            return 1.0 if set(rule['names']).issubset(set(bookmark_bar_folders_names)) else 0.0
        except (KeyError, TypeError, AttributeError):
            return 0.0
    return 0.0

def check_vim_hlsearch_enabled__5438ce42ea45fa77c023ecd730e398a5(result, expected, **options):
    """Check if the vim hlsearch configuration output matches expected.

    Args:
        result: Output from vim config check command
        expected: Expected configuration strings (dict with include/exclude lists)
        **options: Additional options

    Returns:
        float: 1.0 if all expected strings are present and no excluded strings, 0.0 otherwise
    """
    if result is None:
        return 0.0
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_liveserver_extension__a342948a(actual: str, rules: Dict, **options):
    """
    Check if Live Server extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def is_extension_installed__9333e0a8(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_table_row_added__fcf0fdbb(docx_file, expected):
    """
    Check if a new row was added to the first table with expected values.

    Args:
        docx_file: Path to the DOCX file
        expected: Dict with 'expected_rows' (int), 'last_row_values' (list of str)

    Returns:
        float: 1.0 if table has expected row count and last row matches, 0.0 otherwise
    """
    if not docx_file:
        return 0.0
    try:
        doc = Document(docx_file)
    except Exception as e:
        logger.error(f'Error opening document: {e}')
        return 0.0
    expected_rows = expected.get('expected_rows', 7)
    last_row_values = expected.get('last_row_values', ['2024', '50000'])
    if not doc.tables:
        logger.error('Document has no tables')
        return 0.0
    table = doc.tables[0]
    if len(table.rows) != expected_rows:
        logger.info(f'Table row count mismatch: expected {expected_rows}, got {len(table.rows)}')
        return 0.0
    if len(table.rows) < 2:
        logger.error('Table has no data rows')
        return 0.0
    last_row = table.rows[-1]
    actual_values = [cell.text.strip() for cell in last_row.cells]
    if len(actual_values) < 2:
        logger.error(f'Last row has fewer than 2 cells: {actual_values}')
        return 0.0
    for (i, expected_val) in enumerate(last_row_values):
        if actual_values[i].strip() != expected_val:
            logger.info(f"Last row cell {i} mismatch: expected '{expected_val}', got '{actual_values[i]}'")
            return 0.0
    if len(table.rows) >= 2:
        second_to_last_row = table.rows[-2]
        second_to_last_values = [cell.text.strip() for cell in second_to_last_row.cells]
        if len(second_to_last_values) < 2 or all((not val for val in second_to_last_values[:2])):
            logger.info('Second-to-last row is empty - appears row was replaced rather than added')
            return 0.0
    return 1.0

def is_expected_search_query__f9cee6e4(active_tab_info: Dict[str, str], rules: Dict[str, Any]) -> float:
    """
    Check if the active tab URL matches the expected search query pattern.
    Variation 5: Search for first name from cell B10 (Vincenza).
    """
    if not active_tab_info:
        return 0.0
    expected = rules['expect']
    pattern = expected['pattern']
    matched = re.search(pattern, active_tab_info['url'])
    if matched:
        return 1.0
    return 0.0

def is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_task_verify_4(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.
    This is a fixed version that uses the configurable folder names from rules['names']
    instead of hardcoded 'Liked Authors'.
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'liked_authors_websites_urls':
        folder_names = rule.get('names', [])
        if not folder_names:
            logger.error("No folder names specified in rules['names']")
            return 0.0
        target_folder = None
        for folder_name in folder_names:
            target_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder' and bookmark['name'] == folder_name), None)
            if target_folder:
                logger.info(f"Found folder: '{folder_name}'")
                break
        if target_folder:
            folder_urls = [bookmark['url'] for bookmark in target_folder['children'] if bookmark['type'] == 'url']
            logger.info(f"Folder '{target_folder['name']}' URLs: {folder_urls}")
            urls = rule['urls']
            for (idx, url) in enumerate(urls):
                if isinstance(url, str):
                    urls[idx] = [url]
            combinations = product(*urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    return 1.0
            return 0.0
        else:
            logger.error(f'None of the specified folders found: {folder_names}')
            return 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_gdrive_eml_files__74b11cf6(result_state: List[str], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if the downloaded .eml files from Google Drive meet the expected criteria.

    This function verifies that:
    1. The expected number of .eml files were downloaded
    2. Each .eml file contains the required headers (indicating proper Thunderbird export)
    3. Each filename matches the email's subject line (as required by the instruction)
    4. The expected subjects are present (if specified in rules)

    Args:
        result_state: List of local file paths to downloaded .eml files
        expected_state: Dict containing rules (when type='rule', this IS the rules dict):
            - 'file_count': Expected number of files
            - 'check_headers': List of required headers to verify in each file
            - 'expected_subjects': (optional) List of expected email subjects
        **options: Additional options

    Returns:
        float: Score (1.0 if all checks pass, 0.0 otherwise)
    """
    expected_file_count = expected_state.get('file_count', 0)
    required_headers = expected_state.get('check_headers', [])
    expected_subjects = expected_state.get('expected_subjects', [])
    if result_state is None:
        return 0.0
    valid_files = [f for f in result_state if f is not None]
    if len(valid_files) != expected_file_count:
        print(f'File count mismatch: expected {expected_file_count}, got {len(valid_files)}')
        return 0.0
    subjects_found = []
    for file_path in valid_files:
        if not os.path.exists(file_path):
            print(f'File does not exist: {file_path}')
            return 0.0
        if not file_path.endswith('.eml'):
            print(f'File is not a .eml file: {file_path}')
            return 0.0
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            for header in required_headers:
                if header not in content:
                    print(f"Required header '{header}' not found in {file_path}")
                    return 0.0
        except Exception as e:
            print(f'Error reading file {file_path}: {e}')
            return 0.0
        subject = extract_subject_from_eml(file_path)
        if not subject:
            print(f'Could not extract subject from {file_path}')
            return 0.0
        subjects_found.append(subject)
        actual_filename = os.path.basename(file_path)[:-4]
        expected_filename = normalize_filename(subject)
        if actual_filename != expected_filename:
            print(f'Filename does not match subject:')
            print(f'  Filename: {actual_filename}')
            print(f'  Subject:  {subject}')
            print(f'  Expected: {expected_filename}')
            return 0.0
    if expected_subjects:
        expected_subjects_normalized = sorted([s.strip() for s in expected_subjects])
        subjects_found_normalized = sorted([s.strip() for s in subjects_found])
        if expected_subjects_normalized != subjects_found_normalized:
            print(f'Subject mismatch:')
            print(f'  Expected: {expected_subjects_normalized}')
            print(f'  Found:    {subjects_found_normalized}')
            return 0.0
    return 1.0

def check_file_exists_webp__82cfb1cc(src_path, rule):
    """
    Check if the WebP file exists and is readable
    Variation for task 3c8f201a-009d-4bbe-8b65-a6f8b35bb57f_task_verify_5
    """
    if src_path is None:
        return 0.0
    try:
        if not os.path.exists(src_path):
            logger.error(f'File does not exist: {src_path}')
            return 0.0
        if not src_path.lower().endswith('.webp'):
            logger.error(f'File is not a WebP file: {src_path}')
            return 0.0
        file_size = os.path.getsize(src_path)
        if file_size == 0:
            logger.error(f'File is empty: {src_path}')
            return 0.0
        logger.debug(f'WebP file exists with size {file_size} bytes')
        return 1.0
    except Exception as e:
        logger.error(f'Error checking WebP file: {e}')
        return 0.0

def check_table_completeness__5067b5ea(result: Dict[str, Any], expected: Any, **options) -> float:
    """
    Check if table has required number of complete rows.

    Args:
        result: Completeness stats from getter with total_rows and complete_rows
        expected: Dict with:
            - min_complete_rows: Minimum number of complete rows required
            - require_all_complete: Whether all rows must be complete (default: True)

    Returns:
        Score between 0.0 and 1.0
    """
    min_complete_rows = expected.get('min_complete_rows', 1)
    require_all_complete = expected.get('require_all_complete', True)
    total_rows = result.get('total_rows', 0)
    complete_rows = result.get('complete_rows', 0)
    if complete_rows < min_complete_rows:
        return 0.0
    if require_all_complete:
        if total_rows > 0 and complete_rows == total_rows:
            return 1.0
        else:
            return 0.0
    else:
        if total_rows == 0:
            return 0.0
        return complete_rows / total_rows

def check_webext_manifest__b9b28e4a(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_csv_table_mech__9629fd7e(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_recreation_html_element__dbcda1af4d2b810312479cfb54a15ab6(result, expected, **options):
    """
    Compare HTML element extraction result with expected value.
    Used for recreation.gov HTML parsing tasks.
    Verifies both table structure and Echo Canyon specific content.

    Args:
        result: Dict from getter function containing:
            - has_table_header: bool - whether the availability table header exists
            - has_echo_canyon_text: bool - whether "Echo Canyon" text appears in page
            - echo_canyon_in_url: bool - whether Echo Canyon is in the URL
            - page_title: str - page title for context
            - url: str - current page URL
        expected: Expected dict structure from rules
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    logger.info(f'[DEBUG] check_recreation_html_element called with result: {result}')
    logger.info(f'[DEBUG] check_recreation_html_element called with expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, got {type(result)}: {result}')
        return 0.0
    has_table_header = result.get('has_table_header', False)
    if not has_table_header:
        logger.info('[DEBUG] Availability table header not found, returning 0.0')
        return 0.0
    has_echo_canyon_text = result.get('has_echo_canyon_text', False)
    echo_canyon_in_url = result.get('echo_canyon_in_url', False)
    if not (has_echo_canyon_text or echo_canyon_in_url):
        logger.info("[DEBUG] 'Echo Canyon' not found in page content or URL, returning 0.0")
        logger.info(f"[DEBUG] Page URL: {result.get('url', 'N/A')}")
        logger.info(f"[DEBUG] Page title: {result.get('page_title', 'N/A')}")
        return 0.0
    logger.info('[DEBUG] All checks passed:')
    logger.info(f'[DEBUG] - Table header found: {has_table_header}')
    logger.info(f'[DEBUG] - Echo Canyon in page text: {has_echo_canyon_text}')
    logger.info(f'[DEBUG] - Echo Canyon in URL: {echo_canyon_in_url}')
    logger.info(f"[DEBUG] - Page URL: {result.get('url', 'N/A')}")
    return 1.0

def check_chrome_setting__df4ecdef(result, expected, **options):
    """
    Check if the Chrome Homepage URL setting contains the expected value.
    Also verifies that restore_on_startup is set correctly to actually open the homepage.

    Args:
        result: Dictionary with 'setting_value' and 'restore_on_startup' keys (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if both conditions met, 0.0 otherwise
    """
    setting_value = result.get('setting_value', [])
    restore_on_startup = result.get('restore_on_startup', 1)
    if expected['type'] == 'contains':
        url_present = False
        if isinstance(setting_value, list):
            url_present = expected['value'] in setting_value
        else:
            url_present = expected['value'] == setting_value
        restore_correct = restore_on_startup in [4, 5]
        logger.info(f'[CHROME_CHECK] URL present: {url_present}, restore_on_startup: {restore_on_startup}, restore_correct: {restore_correct}')
        if url_present and restore_correct:
            return 1.0
        else:
            return 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_extension_uppercase__de1c34d5(result, expected, **options):
    """Check if all image files have uppercase extension.

    The evaluator runs an inline shell command that:
    1. Navigates to /home/user/Pictures/
    2. Computes SHA256 hashes for all *.JPG and *.jpg files
    3. Outputs JSON mapping: {hash: filename} e.g. {"abc123...": "picture1.JPG"}

    This function verifies that all expected image hashes (representing the 3 downloaded
    images) map to filenames ending with the expected extension (.JPG).

    Args:
        result: JSON string/dict mapping image SHA256 hashes to filenames
        expected: Rules dict with:
            - expected_extension: Target extension (e.g., ".JPG")
            - image_hashes: List of SHA256 hashes for the images that should be renamed
        **options: Additional options

    Returns:
        float: 1.0 if all image hashes exist and have correct extension, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if isinstance(result, str):
        result = result.strip().replace("'", '"')
        try:
            result = json.loads(result)
        except:
            return 0.0
    expected_ext = expected.get('expected_extension', '')
    image_hashes = expected.get('image_hashes', [])
    for img_hash in image_hashes:
        if img_hash not in result:
            return 0.0
        filename = result[img_hash]
        if not filename.endswith(expected_ext):
            return 0.0
    return 1.0

def check_extension_enabled__6f71517e0c42749af1a6363d9f36e224(result, expected, **options):
    """
    Check if the extension is installed and enabled.

    Args:
        result: Dict from getter with 'path' and 'enabled' keys
        expected: Dict with 'extension_path' and 'should_be_enabled' keys
        **options: Additional options (unused)

    Returns:
        float: 1.0 if extension is installed at correct path and has correct enabled state, 0.0 otherwise
    """
    if not result:
        return 0.0
    expected_path = expected.get('extension_path', '')
    should_be_enabled = expected.get('should_be_enabled', True)
    if result.get('path') == expected_path and result.get('enabled') == should_be_enabled:
        return 1.0
    return 0.0

def check_recreation_html_element__fa1e76c31141f93d38de11c4bb8239cf(result, expected, **options):
    """
    Compare HTML element extraction result with expected value.
    Used for recreation.gov Devil's Garden search verification.

    This metric verifies:
    1. The correct location (Devil's Garden) is being viewed
    2. A reservation table is present
    3. Availability data is displayed
    4. Reservation dates are sorted (earliest first) for "soonest" search
    5. The earliest reservation is identified/displayed

    Args:
        result: Dict from getter function containing:
                - location_verified: bool (Devil's Garden found in URL/content)
                - reservation_table_present: bool (table headers found)
                - has_availability_data: bool (reservation data displayed)
                - dates_sorted: bool (dates in ascending order)
                - earliest_reservation_identified: bool (earliest date is visible/highlighted)
                - reservation_dates: list (extracted reservation dates)
                - url: str (current page URL)
                - page_title: str (page title)
        expected: Expected dict structure from rules containing:
                - location: str (expected location name)
                - location_verified: bool (expected verification state)
                - reservation_table_present: bool (expected table presence)
                - has_availability_data: bool (expected data presence)
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    logger.info(f'[DEBUG] check_recreation_html_element called with result: {result}')
    logger.info(f'[DEBUG] check_recreation_html_element called with expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, got {type(result)}: {result}')
        return 0.0
    expected_location = expected.get('location', "Devil's Garden")
    expected_location_verified = expected.get('location_verified', True)
    expected_table_present = expected.get('reservation_table_present', True)
    expected_has_data = expected.get('has_availability_data', True)
    logger.info(f'[DEBUG] Expected values: location={expected_location}, location_verified={expected_location_verified}, table_present={expected_table_present}, has_data={expected_has_data}')
    location_verified = result.get('location_verified', False)
    if expected_location_verified and (not location_verified):
        logger.info(f'[DEBUG] ✗ Location not verified - {expected_location} not found in URL or page content')
        return 0.0
    logger.info(f'[DEBUG] ✓ Location verified: {expected_location} confirmed')
    reservation_table_present = result.get('reservation_table_present', False)
    if expected_table_present and (not reservation_table_present):
        logger.info('[DEBUG] ✗ Reservation table not present - camp-sortable-column-header elements not found')
        return 0.0
    logger.info('[DEBUG] ✓ Reservation table present')
    has_availability_data = result.get('has_availability_data', False)
    if expected_has_data and (not has_availability_data):
        logger.info('[DEBUG] ✗ Availability data not found')
        return 0.0
    logger.info('[DEBUG] ✓ Availability data found')
    dates_sorted = result.get('dates_sorted', False)
    reservation_dates = result.get('reservation_dates', [])
    if len(reservation_dates) == 0:
        logger.info('[DEBUG] ✗ No reservation dates found - search may be incomplete')
        return 0.5
    if not dates_sorted:
        logger.info(f'[DEBUG] ✗ Reservation dates not sorted in ascending order')
        logger.info(f'[DEBUG]   Found dates: {reservation_dates}')
        return 0.6
    logger.info(f'[DEBUG] ✓ Reservation dates are sorted (earliest first)')
    logger.info(f'[DEBUG]   Dates: {(reservation_dates[:5] if len(reservation_dates) > 5 else reservation_dates)}')
    earliest_reservation_identified = result.get('earliest_reservation_identified', False)
    if not earliest_reservation_identified:
        logger.info('[DEBUG] ✗ Earliest reservation not clearly identified')
        return 0.8
    logger.info(f'[DEBUG] ✓ Earliest reservation identified')
    if len(reservation_dates) > 0:
        logger.info(f'[DEBUG]   Earliest date: {reservation_dates[0]}')
    score = 1.0
    logger.info(f'[DEBUG] ========================================')
    logger.info(f'[DEBUG] All checks passed - Final score: {score}')
    logger.info(f'[DEBUG] ========================================')
    logger.info(f"[DEBUG] Page URL: {result.get('url', 'N/A')}")
    logger.info(f"[DEBUG] Page title: {result.get('page_title', 'N/A')}")
    return score

def is_extension_installed__1d58e082(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_recreation_html_element__ed948b8fac72e2c98dd5028aabd38bf3(result, expected, **options):
    """
    Verify that the user successfully located the closest available campsite for Cedar Breaks.

    This checks:
    1. Cedar Breaks was searched for (appears in page content)
    2. Campsite results were found and displayed
    3. The closest campsite was identified (first result or sorted by distance)

    Args:
        result: Dict from getter function containing:
            - cedar_breaks_found: bool indicating if Cedar Breaks is in the page
            - has_campsite_results: bool indicating if campsite data is shown
            - campsites: list of campsite data
            - closest_identified: bool indicating if closest was determined
        expected: Expected dict structure from rules
        **options: Additional comparison options

    Returns:
        float: Score between 0.0 (fail) and 1.0 (pass)
    """
    logger.info(f'[DEBUG] check_recreation_html_element called with result: {result}')
    logger.info(f'[DEBUG] check_recreation_html_element called with expected: {expected}')
    if result is None:
        logger.info('[DEBUG] Result is None, returning 0.0')
        return 0.0
    if not isinstance(result, dict):
        logger.error(f'[DEBUG] Result is not a dict, got {type(result)}: {result}')
        return 0.0
    cedar_breaks_found = result.get('cedar_breaks_found', False)
    if not cedar_breaks_found:
        logger.info('[DEBUG] Cedar Breaks not found in page content - task incomplete')
        return 0.0
    has_results = result.get('has_campsite_results', False)
    if not has_results:
        logger.info('[DEBUG] No campsite results found - task incomplete')
        return 0.0
    closest_identified = result.get('closest_identified', False)
    if not closest_identified:
        logger.info('[DEBUG] Closest available campsite not identified - task incomplete')
        return 0.0
    campsites = result.get('campsites', [])
    if len(campsites) == 0:
        logger.info('[DEBUG] No campsites found in results - task incomplete')
        return 0.0
    available_campsites = [c for c in campsites if c.get('available', False)]
    if len(available_campsites) == 0:
        logger.info('[DEBUG] No available campsites found - task incomplete')
        return 0.0
    logger.info(f'[DEBUG] Task completed successfully! Cedar Breaks: {cedar_breaks_found}, Results: {has_results}, Closest available: {closest_identified}, Available campsites: {len(available_campsites)}/{len(campsites)}')
    logger.info(f"[DEBUG] Closest available campsite: {available_campsites[0].get('name', 'Unknown')}")
    return 1.0

def check_restclient_extension__d9713508(actual: str, rules: Dict, **options):
    """
    Check if REST Client extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_gdrive_pdf__f93e2dd2(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_table_sorted_with_row_integrity__b82f9c9e(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if a table is sorted by one column while preserving row integrity.

    Args:
        result: Dict with 'columns' containing column data
        expected: Dict with 'sort_column', 'order', 'original_pairs' (known Order ID -> Mark mappings)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    columns = result.get('columns', {})
    sort_column = expected.get('sort_column', 'C')
    order = expected.get('order', 'desc')
    original_pairs = expected.get('original_pairs', {})
    if sort_column not in columns:
        return 0.0
    marks = columns.get('C', [])
    order_ids = columns.get('B', [])
    if len(marks) == 0 or len(order_ids) == 0:
        return 0.0
    if len(marks) != len(order_ids):
        return 0.0
    numeric_marks = [v for v in marks if isinstance(v, (int, float))]
    if order == 'desc':
        expected_sorted = sorted(numeric_marks, reverse=True)
    else:
        expected_sorted = sorted(numeric_marks)
    sorted_score = sum((1 for (i, v) in enumerate(numeric_marks) if i < len(expected_sorted) and v == expected_sorted[i])) / len(numeric_marks) if numeric_marks else 0.0
    current_pairs = {}
    for (i, (order_id, mark)) in enumerate(zip(order_ids, marks)):
        if isinstance(mark, (int, float)):
            current_pairs[order_id] = mark
    integrity_score = 0.0
    if original_pairs and current_pairs:
        matches = 0
        total = 0
        for (order_id, original_mark) in original_pairs.items():
            total += 1
            if order_id in current_pairs and current_pairs[order_id] == original_mark:
                matches += 1
        integrity_score = matches / total if total > 0 else 0.0
    else:
        integrity_score = 1.0
    return (sorted_score + integrity_score) / 2.0

def check_table_row_count_by_index__f500041f(result, expected, **options):
    """Check table row count and header row content."""
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    result_index = result.get('table_index', -1)
    expected_index = expected.get('table_index', -1)
    if result_index != expected_index:
        return 0.0
    result_rows = result.get('row_count', 0)
    min_rows = expected.get('min_rows', 0)
    if result_rows < min_rows:
        return 0.0
    first_row_cells = result.get('first_row_cells', [])
    expected_graphemes = ['ai', 'ee', 'igh', 'ow', 'oo']
    normalized_cells = [cell.lower().strip() for cell in first_row_cells]
    if len(normalized_cells) < len(expected_graphemes):
        return 0.0
    for (i, expected_grapheme) in enumerate(expected_graphemes):
        if i >= len(normalized_cells):
            return 0.0
        if expected_grapheme not in normalized_cells[i]:
            return 0.0
    return 1.0

def is_expected_bookmarks__ee2a139f7ed95e764d31e544c144187a(result_state, expected_state, **options):
    """
    Custom metric that checks for bookmarks in a folder with a configurable name.
    This version reads the folder name from expected_state['rules']['names'][0]
    instead of hardcoding 'Liked Authors'.

    Args:
        result_state: The current bookmarks structure from get_bookmarks getter
        expected_state: Expected state with rules including:
            - type: 'bookmark_folder_websites_urls' or 'liked_authors_websites_urls'
            - names: List with folder name as first element
            - urls: List of lists of possible URLs for each bookmark
        **options: Additional options

    Returns:
        float: 1.0 if bookmarks match, 0.0 otherwise
    """
    bookmarks = result_state
    rule = expected_state['rules']
    if rule['type'] in ['bookmark_folder_websites_urls', 'liked_authors_websites_urls']:
        folder_name = rule['names'][0] if rule.get('names') else None
        if not folder_name:
            logger.error("No folder name specified in rules['names']")
            return 0.0
        target_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder' and bookmark['name'] == folder_name), None)
        if target_folder:
            logger.info(f"'{folder_name}' folder exists")
            folder_urls = [bookmark['url'] for bookmark in target_folder['children'] if bookmark['type'] == 'url']
            logger.info(f"Here is the '{folder_name}' folder's urls: {folder_urls}")
            urls = rule['urls']
            for (idx, url) in enumerate(urls):
                if isinstance(url, str):
                    urls[idx] = [url]
            combinations = product(*urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    return 1.0
            return 0.0
        else:
            logger.info(f"'{folder_name}' folder does not exist")
            return 0.0
    return 0.0

def is_extension_installed__9024492c(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_chrome_experiments_contains_any__70b32da51f968d0aa45e836eecc52bdb(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the enabled Chrome experiments contain ANY of the expected experiments.

    Args:
        result: List of enabled experiment names
        expected: Dict with 'experiment_names' key containing a list of potential experiments
        **options: Additional options

    Returns:
        1.0 if at least one expected experiment is in the list, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    expected_experiments = expected.get('experiment_names', [])
    if not expected_experiments:
        logger.error('No experiment_names specified in expected rules')
        return 0.0
    if not isinstance(expected_experiments, list):
        logger.error(f'experiment_names is not a list: {type(expected_experiments)}')
        return 0.0
    found_experiments = [exp for exp in expected_experiments if exp in result]
    if found_experiments:
        logger.info(f'Found at least one expected experiment: {found_experiments}')
        return 1.0
    else:
        logger.info(f'No expected experiments found. Expected any of: {expected_experiments}, Got: {result}')
        return 0.0

def check_contains_url__a2f501d6(result, expected, **options):
    """
    Check if result contains the expected URL.

    Args:
        result: List of URLs from getter
        expected: Dict with 'url' key containing expected URL
        **options: Additional options

    Returns:
        float: 1.0 if URL is found, 0.0 otherwise
    """
    if not isinstance(result, list):
        return 0.0
    expected_url = expected.get('url', '')
    if not expected_url:
        return 0.0
    for url in result:
        if expected_url.lower() in url.lower():
            return 1.0
    return 0.0

def check_extension_version__a366045b(result, rules) -> float:
    """Check if extension version matches expected value.

    Args:
        result: Version string from getter
        rules: Dict with 'expected' key

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_version = rules.get('expected')
    if result == expected_version:
        return 1.0
    return 0.0

def check_gdrive_file_exists__b791b56367b138c183fb013d8a0662b9(result: bool, expected: Dict[str, Any], **options) -> float:
    """Check if file exists in Google Drive folder.

    Args:
        result: Boolean from getter indicating if file exists
        expected: Expected rules dict (currently unused, but required for framework)
        **options: Additional options

    Returns:
        float: 1.0 if file exists, 0.0 otherwise
    """
    if result:
        logger.info('File exists in Google Drive folder as expected')
        return 1.0
    else:
        logger.info('File not found in Google Drive folder')
        return 0.0

def check_webext_manifest__81e9cd46(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    expected_has_empty_description = expected.get('has_empty_description', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    if expected_has_empty_description:
        checks += 1
        description = manifest.get('description', '')
        if not description or description == '':
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_extension_name_path__d51a0e10(result, expected, **options):
    """
    Check if an extension with the specified name and path exists in Chrome.

    Args:
        result: List of dicts with extension info [{"name": "...", "path": "..."}]
        expected: Dict with expected extension info {"name": "...", "path": "..."}
        **options: Additional options

    Returns:
        float: 1.0 if extension with matching name and path is found, 0.0 otherwise
    """
    if not result:
        logger.error('No extensions found')
        return 0.0
    if not expected:
        logger.error('No expected extension data provided')
        return 0.0
    expected_name = expected.get('name', '')
    expected_path = expected.get('path', '')
    if not expected_name or not expected_path:
        logger.error(f"Expected name or path not specified: name='{expected_name}', path='{expected_path}'")
        return 0.0
    logger.info(f"Looking for extension with name='{expected_name}' and path='{expected_path}'")
    logger.info(f'Found {len(result)} installed extensions')
    for extension in result:
        actual_name = extension.get('name', '')
        actual_path = extension.get('path', '')
        logger.debug(f"Checking extension: name='{actual_name}', path='{actual_path}'")
        if actual_name == expected_name and actual_path == expected_path:
            logger.info(f"✓ Found matching extension: name='{actual_name}', path='{actual_path}'")
            return 1.0
    logger.error(f"✗ Extension with name='{expected_name}' and path='{expected_path}' not found")
    logger.error(f'Installed extensions:')
    for ext in result:
        logger.error(f"  - name='{ext.get('name', '')}', path='{ext.get('path', '')}'")
    return 0.0

def check_chrome_experiments_contains_all__e998f78abb27064086318477b860256b(result: List[str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the enabled Chrome experiments contain ALL of the expected experiments.

    Args:
        result: List of enabled experiment names
        expected: Dict with 'experiment_names' key containing a list of required experiments
        **options: Additional options

    Returns:
        1.0 if all expected experiments are in the list, 0.0 otherwise
    """
    if not isinstance(result, list):
        logger.error(f'Result is not a list: {type(result)}')
        return 0.0
    expected_experiments = expected.get('experiment_names', [])
    if not expected_experiments:
        logger.error('No experiment_names specified in expected rules')
        return 0.0
    if not isinstance(expected_experiments, list):
        logger.error(f'experiment_names is not a list: {type(expected_experiments)}')
        return 0.0
    missing_experiments = [exp for exp in expected_experiments if exp not in result]
    if not missing_experiments:
        logger.info(f'All expected experiments found: {expected_experiments}')
        return 1.0
    else:
        logger.info(f'Missing experiments: {missing_experiments}. Enabled: {result}')
        return 0.0

def is_expected_search_query__8164e914(active_tab_info: Dict[str, str], rules: Dict[str, Any]) -> float:
    """
    Check if the active tab URL matches the expected search query pattern.
    Variation 2: Search for ID number from cell H6 (2468).
    """
    if not active_tab_info:
        return 0.0
    expected = rules['expect']
    pattern = expected['pattern']
    matched = re.search(pattern, active_tab_info['url'])
    if matched:
        return 1.0
    return 0.0

def check_gdrive_files__f478ac41(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_csv_table_bio__a4ea4d07(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    total_checks = 0
    if expected_headers:
        total_checks += 1
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            total_checks += 1
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_csv_table_cs3y__26538d5a(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def is_expected_bookmarks__2ad9387a_65d8_4e33_ad5b_7580065a27ca_aug_1_task_verify_1(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.
    This version uses subset checking for bookmark_bar_folders_names to verify
    that the expected folders are PRESENT, not that they are the ONLY folders.

    Args:
        bookmarks: Bookmarks data from Chrome
        rule: Rule configuration with type and expected values

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if 'bookmark_bar' not in bookmarks:
        return 0.0
    bookmark_bar = bookmarks['bookmark_bar']
    if not bookmark_bar or 'children' not in bookmark_bar:
        return 0.0
    children = bookmark_bar['children']
    if children is None:
        children = []
    if rule['type'] == 'bookmark_bar_folders_names':
        bookmark_bar_folders_names = [bookmark['name'] for bookmark in children if bookmark.get('type') == 'folder' and 'name' in bookmark]
        expected_names = rule.get('names', [])
        return 1.0 if set(expected_names).issubset(set(bookmark_bar_folders_names)) else 0.0
    elif rule['type'] == 'bookmark_bar_websites_urls':
        bookmark_bar_websites_urls = [bookmark['url'] for bookmark in children if bookmark.get('type') == 'url' and 'url' in bookmark]
        return 1.0 if set(bookmark_bar_websites_urls) == set(rule.get('urls', [])) else 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_html_contains__94d6a0b8(result, rules) -> float:
    """Check if HTML content contains expected text.

    Args:
        result: HTML content string from getter
        rules: Dict with 'expected_text' key

    Returns:
        float: 1.0 if contains expected text, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expected_text = rules.get('expected_text', '')
    if expected_text in result:
        return 1.0
    return 0.0

def check_websites_valid__af875e9f(result, expected, **options):
    """Check if websites are filled and look like valid URLs.

    Args:
        result: List of websites from Excel
        expected: Dict (can be empty or contain validation params)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    valid_count = 0
    for website in result:
        if website and website.strip():
            website_lower = website.lower().strip()
            if any([website_lower.startswith('http://'), website_lower.startswith('https://'), website_lower.startswith('www.'), '.' in website_lower]):
                valid_count += 1
    return valid_count / len(result)

def check_gdrive_files__7fbebc15(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_chrome_hardware_accel__8f79fa7c(result, expected, **options):
    """Check if Chrome hardware acceleration is in expected state.

    Args:
        result: Hardware acceleration settings dict from getter
        expected: Expected rules dict with 'enabled' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_enabled = expected.get('enabled', True)
    actual_enabled = result.get('enabled', True)
    logger.info(f'Expected hardware acceleration enabled: {expected_enabled}')
    logger.info(f'Actual hardware acceleration enabled: {actual_enabled}')
    return 1.0 if actual_enabled == expected_enabled else 0.0

def check_csv_export_and_chrome__4e1aca7b(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if CSV file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with csv_path and csv_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f"CSV file exists at {result.get('expected_file_path')}: +0.5")
    else:
        logger.warning(f"CSV file does not exist at {result.get('expected_file_path')}")
    chrome_tabs = result.get('chrome_tabs', [])
    csv_url_pattern = expected.get('csv_url_pattern', '')
    csv_opened = False
    for tab_url in chrome_tabs:
        if csv_url_pattern and csv_url_pattern in tab_url:
            csv_opened = True
            logger.info(f'Found exact CSV URL pattern in tab: {tab_url}')
            break
        elif tab_url.endswith('.csv') and 'annual-enterprise-survey-2021-financial-year-provisional' in tab_url:
            csv_opened = True
            logger.info(f'Found CSV file in tab: {tab_url}')
            break
    if csv_opened:
        score += 0.5
        logger.info(f'CSV file opened in Chrome: +0.5')
    else:
        logger.warning(f'CSV file not found in Chrome tabs. Expected pattern: {csv_url_pattern}, Found tabs: {chrome_tabs}')
    logger.info(f"Final score: {score} (file_exists: {result.get('file_exists')}, csv_opened: {csv_opened})")
    return score

def check_table_count_only__db5e4b5e(result, expected, **options):
    """
    Check table count and bulleted list requirements.

    Verifies:
    1. Table count decreased from 13 to 12 (one table removed)
    2. A bulleted list was created
    3. The list contains exactly 7 items with the expected consonants: p, b, t, d, k, c, g
    4. The bulleted list appears at an appropriate position (near where second table was)
    5. CRITICAL: The specific table containing the consonants was removed (not just any table)
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        logger.error('Invalid result or expected format')
        return 0.0
    result_count = result.get('table_count', 0)
    expected_count = expected.get('table_count', 0)
    table_count_correct = result_count == expected_count
    has_bulleted_list = result.get('has_bulleted_list', False)
    expected_has_list = expected.get('has_bulleted_list', True)
    bulleted_list_exists = has_bulleted_list == expected_has_list
    bulleted_list_items = result.get('bulleted_list_items', [])
    expected_consonants = set(expected.get('expected_consonants', ['p', 'b', 't', 'd', 'k', 'c', 'g']))
    found_items = set()
    for item in bulleted_list_items:
        item_lower = item.strip().lower()
        for consonant in expected_consonants:
            if consonant in item_lower:
                found_items.add(consonant)
    consonants_found = expected_consonants.issubset(found_items)
    correct_item_count = len(bulleted_list_items) == 7
    if not correct_item_count:
        logger.info(f'✗ Expected exactly 7 bulleted items, found {len(bulleted_list_items)}')
    list_location_index = result.get('list_location_index', -1)
    table_positions = result.get('table_positions', [])
    positional_check_passed = False
    if list_location_index >= 0 and len(table_positions) >= 2:
        if list_location_index < table_positions[min(2, len(table_positions) - 1)]:
            positional_check_passed = True
            logger.info(f'✓ List positioned appropriately at index {list_location_index}, before table at {table_positions[min(2, len(table_positions) - 1)]}')
        else:
            logger.info(f'✗ List position {list_location_index} seems late, expected before table at {table_positions[min(2, len(table_positions) - 1)]}')
    elif list_location_index >= 0:
        positional_check_passed = True
        logger.info(f'✓ List found at position {list_location_index}')
    else:
        logger.info(f'✗ List location not found')
    tables_with_consonants = result.get('tables_with_consonants', [])
    consonant_table_removed = len(tables_with_consonants) == 0
    if not consonant_table_removed:
        logger.info(f'✗ CRITICAL: Found {len(tables_with_consonants)} table(s) still containing consonants at indices {tables_with_consonants}')
        logger.info(f'   This indicates the second table with consonants was NOT removed, or wrong table was removed')
    else:
        logger.info(f'✓ CRITICAL: No tables containing the expected consonants found - specific table was removed')
    checks_passed = 0
    total_checks = 6
    if table_count_correct:
        checks_passed += 1
        logger.info(f'✓ Table count matches: {result_count}')
    else:
        logger.info(f'✗ Table count mismatch: got {result_count}, expected {expected_count}')
    if bulleted_list_exists:
        checks_passed += 1
        logger.info(f'✓ Bulleted list exists')
    else:
        logger.info(f'✗ No bulleted list found')
    if consonants_found:
        checks_passed += 1
        logger.info(f'✓ All expected consonants found: {found_items}')
    else:
        missing = expected_consonants - found_items
        logger.info(f'✗ Missing consonants: {missing}, found: {found_items}')
    if correct_item_count:
        checks_passed += 1
        logger.info(f'✓ Correct number of list items (7)')
    else:
        logger.info(f'✗ Expected 7 list items, found {len(bulleted_list_items)}')
    if positional_check_passed:
        checks_passed += 1
    if consonant_table_removed:
        checks_passed += 1
        logger.info(f'✓ Specific consonant table was removed')
    else:
        logger.info(f'✗ Consonant table still exists or wrong table removed')
    score = checks_passed / total_checks
    logger.info(f'Final score: {score:.2f} ({checks_passed}/{total_checks} checks passed)')
    return score

def check_extension_description__db4e5321(result, rules) -> float:
    """Check if extension description matches expected value.

    Args:
        result: Description string from getter
        rules: Dict with 'expected' key

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_description = rules.get('expected')
    if result == expected_description:
        return 1.0
    return 0.0

def check_table_centered__bbe9b961(result, expected, **options):
    """Check if table is vertically centered.

    Args:
        result: Table position dict from getter
        expected: Dict with 'target_top' and 'tolerance'
        **options: Additional options

    Returns:
        float: 1.0 if within tolerance, 0.0 otherwise
    """
    if result is None:
        return 0.0
    target_top = expected.get('target_top', 0)
    tolerance = expected.get('tolerance', 0)
    if abs(result['top'] - target_top) <= tolerance:
        return 1.0
    else:
        return 0.0

def check_account_server_url__9e272861(result: str, expected: Dict[str, List[Dict[str, str]]]) -> float:
    """
    Check if account with specific server URL exists in Thunderbird.

    Args:
        result: path to csv file containing account data
        expected: dict with "expect" key containing list of records with "url" field

    Returns:
        float: 1.0 if all expected server URLs are found, 0.0 otherwise
    """
    if result is None:
        return 0.0
    expect_list = expected.get('expect', [])
    if not expect_list:
        return 1.0
    expect_metrics = [False] * len(expect_list)
    try:
        with open(result) as f:
            reader = csv.DictReader(f)
            for rcd in reader:
                for (i, expected_rec) in enumerate(expect_list):
                    if all((rcd.get(k) == v for (k, v) in expected_rec.items())):
                        expect_metrics[i] = True
    except Exception as e:
        logger.error(f'Error reading CSV: {e}')
        return 0.0
    return float(all(expect_metrics))

def check_extension_version__407be0458b7b234fb5401d66a10f5221(result, expected, **options):
    """Check if the extension version matches the expected version.

    Args:
        result: Version string from getter
        expected: Dict with 'version' key containing the expected version
        **options: Additional options (not used)

    Returns:
        float: 1.0 if versions match, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    expected_version = expected.get('version', '')
    if not expected_version:
        return 0.0
    if result == expected_version:
        return 1.0
    else:
        return 0.0

def check_csv_table_cs4y__5116177c(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_extension_enabled__ae6416e4(result: Dict[str, bool], expected: Dict[str, Any], **options) -> float:
    """Check if a specific extension is enabled.

    Args:
        result: Dict mapping extension names to enabled status
        expected: Rules dict with 'extension_name' key
        **options: Additional options

    Returns:
        1.0 if extension is enabled, 0.0 otherwise
    """
    extension_name = expected.get('extension_name', '')
    if not extension_name:
        return 0.0
    if extension_name in result and result[extension_name]:
        return 1.0
    else:
        return 0.0

def is_expected_bookmarks__a82b78bb_aug14(result_state: Dict[str, Any], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if bookmarks match expected state with configurable folder name.

    This metric verifies that a bookmark folder with a specified name exists
    and contains the expected combination of URLs.

    Args:
        result_state: Actual bookmark state from getter
        expected_state: Expected state with rules configuration
        **options: Additional options

    Returns:
        float: Score (1.0 if match, 0.0 otherwise)
    """
    bookmarks = result_state
    rule = expected_state.get('rules', {})
    if rule.get('type') == 'bookmark_folder_urls':
        folder_name = rule.get('names', [None])[0] if rule.get('names') else None
        if not folder_name:
            logger.error("No folder name specified in rules['names']")
            return 0.0
        if 'bookmark_bar' not in bookmarks or 'children' not in bookmarks['bookmark_bar']:
            logger.error('bookmark_bar or its children not found in bookmarks')
            return 0.0
        target_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark.get('type') == 'folder' and bookmark.get('name') == folder_name), None)
        if target_folder:
            logger.info(f"'{folder_name}' folder exists")
            folder_urls = [bookmark['url'] for bookmark in target_folder.get('children', []) if bookmark.get('type') == 'url' and 'url' in bookmark]
            logger.info(f"'{folder_name}' folder's urls: {folder_urls}")
            urls = rule.get('urls', [])
            normalized_urls = []
            for url in urls:
                if isinstance(url, str):
                    normalized_urls.append([url])
                elif isinstance(url, list):
                    normalized_urls.append(url)
                else:
                    logger.error(f'Invalid URL format: {url}')
                    return 0.0
            combinations = product(*normalized_urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    logger.info(f'Found matching combination: {combination}')
                    return 1.0
            logger.info(f'No matching combination found. Expected one of the combinations from {normalized_urls}, got {folder_urls}')
            return 0.0
        else:
            logger.info(f"'{folder_name}' folder does not exist")
            return 0.0
    else:
        raise TypeError(f"{rule.get('type')} not support yet!")

def is_expected_bookmarks__35253b65(bookmarks: Dict[str, Any], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks with both URL and title are in Chrome's bookmark bar.

    This custom metric supports the 'bookmark_bar_websites_with_titles' validation type
    which checks both the URL and title of bookmarks.

    Args:
        bookmarks: Bookmarks data structure from Chrome
        rule: Rule containing validation type and expected bookmarks

    Returns:
        float: 1.0 if all expected bookmarks (URL + title) are present, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] != 'bookmark_bar_websites_with_titles':
        raise TypeError(f"{rule['type']} not supported by this metric!")
    if 'bookmark_bar' not in bookmarks or not bookmarks['bookmark_bar']:
        return 0.0
    if 'children' not in bookmarks['bookmark_bar'] or not bookmarks['bookmark_bar']['children']:
        return 0.0
    bookmark_bar_websites = [{'url': bookmark.get('url', ''), 'title': bookmark.get('name', '')} for bookmark in bookmarks['bookmark_bar']['children'] if bookmark.get('type') == 'url']
    expected_bookmarks = rule['bookmarks']
    for expected in expected_bookmarks:
        expected_url = expected['url']
        expected_title = expected['title']
        found = False
        for actual in bookmark_bar_websites:
            if actual['url'] == expected_url and actual['title'] == expected_title:
                found = True
                break
        if not found:
            return 0.0
    return 1.0

def check_gdrive_file__6bd409d3(result: dict, expected: dict, **options) -> float:
    """
    Check if file exists on Google Drive and meets size requirements.

    Args:
        result: Metadata dict from getter
        expected: Dict with 'min_size' requirement
        **options: Additional options

    Returns:
        float: 1.0 if file exists and meets requirements, 0.0 otherwise
    """
    if result is None or not result.get('exists', False):
        logger.warning('File does not exist on Google Drive')
        return 0.0
    min_size = expected.get('min_size', 0)
    actual_size = result.get('size', 0)
    if actual_size < min_size:
        logger.warning(f'File too small: {actual_size} < {min_size}')
        return 0.0
    logger.info(f'File exists on Google Drive with size {actual_size} bytes')
    return 1.0

def check_csv_file_and_tab__29a47154(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if CSV file exists and is opened in Chrome.

    Args:
        result: Dict with file_exists and chrome_tabs from getter
        expected: Dict with csv_path, expected_urls, and csv_url_pattern
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    score = 0.0
    if result.get('file_exists', False):
        score += 0.5
        logger.info(f'CSV file exists: +0.5')
    else:
        logger.warning(f'CSV file does not exist')
    chrome_tabs = result.get('chrome_tabs', [])
    csv_url_pattern = expected.get('csv_url_pattern', '')
    csv_opened = False
    for tab_url in chrome_tabs:
        if csv_url_pattern in tab_url or tab_url.endswith('.csv'):
            csv_opened = True
            break
    if csv_opened:
        score += 0.5
        logger.info(f'CSV file opened in Chrome: +0.5')
    else:
        logger.warning(f'CSV file not found in Chrome tabs. Tabs: {chrome_tabs}')
    logger.info(f'Final score: {score}')
    return score

def check_docker_extension__acbefae0(actual: str, rules: Dict, **options):
    """
    Check if Docker extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_extension_manifest__b5fa3477caae1bab3fd0d8b19ef4dfc7(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate browser extension manifest.json structure.

    Args:
        result: Manifest data from getter
        expected: Rules dict containing expected values for name, version, background, browser_action
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct field)
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dictionary')
        return 0.0
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    check_background = expected.get('check_background', False)
    check_browser_action = expected.get('check_browser_action', False)
    score = 0.0
    total_checks = 0
    total_checks += 1
    if result.get('name') == expected_name:
        score += 0.25
        logger.info(f'Name matches: {expected_name}')
    else:
        logger.info(f"Name mismatch: expected={expected_name}, got={result.get('name')}")
    total_checks += 1
    if result.get('version') == expected_version:
        score += 0.25
        logger.info(f'Version matches: {expected_version}')
    else:
        logger.info(f"Version mismatch: expected={expected_version}, got={result.get('version')}")
    if check_background:
        total_checks += 1
        background = result.get('background', {})
        if isinstance(background, dict) and 'scripts' in background:
            scripts = background.get('scripts', [])
            if isinstance(scripts, list) and len(scripts) > 0:
                score += 0.25
                logger.info(f'Background scripts found: {scripts}')
            else:
                logger.info('Background scripts field exists but is empty')
        else:
            logger.info('Background section missing or invalid')
    if check_browser_action:
        total_checks += 1
        browser_action = result.get('browser_action', {})
        if isinstance(browser_action, dict) and len(browser_action) > 0:
            has_fields = any((k in browser_action for k in ['default_popup', 'default_icon', 'default_title']))
            if has_fields:
                score += 0.25
                logger.info(f'Browser action found with fields: {list(browser_action.keys())}')
            else:
                logger.info('Browser action exists but missing expected fields')
        else:
            logger.info('Browser action section missing or invalid')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_url_contains__9fa2717e(result: Optional[str], expected: dict, **options) -> float:
    """Check if URL contains expected substring and validates it's a Vicuna paper from 2023.

    Args:
        result: Actual URL string from getter
        expected: Expected rules dict with 'substring', 'keyword', and 'year' values
        **options: Additional comparison options

    Returns:
        1.0 if URL contains substring and meets all requirements, 0.0 otherwise
    """
    if result is None:
        return 0.0
    result_lower = result.lower()
    substring = expected.get('substring', '').lower()
    if not substring:
        return 0.0
    if substring not in result_lower:
        return 0.0
    keyword = expected.get('keyword', 'vicuna').lower()
    if keyword not in result_lower:
        return 0.0
    year_pattern = expected.get('year_pattern', '23(0[1-9]|1[0-2])')
    if re.search(year_pattern, result):
        return 1.0
    if '2023' in result:
        return 1.0
    return 0.0

def is_expected_search_query__df8126bf(active_tab_info: Dict[str, str], rules: Dict[str, Any]) -> float:
    """
    Check if the active tab URL matches the expected search query pattern.
    Variation 7: Search for last name from cell C2 (Abril).
    """
    if not active_tab_info:
        return 0.0
    expected = rules['expect']
    pattern = expected['pattern']
    matched = re.search(pattern, active_tab_info['url'])
    if matched:
        return 1.0
    return 0.0

def check_search_engine__d2ec4a7b(result_state, expected_state, **options):
    """
    Check if the default search engine matches the expected value.

    Args:
        result_state: String representing the current default search engine name
        expected_state: Dict containing the rules (e.g., {'name': 'DuckDuckGo'})
                       When evaluator uses type='rule', the framework passes config['rules'] directly
        **options: Additional options

    Returns:
        float: 1.0 if search engine matches expected, 0.0 otherwise
    """
    if not result_state:
        logger.error('No search engine result returned')
        return 0.0
    expected_name = expected_state.get('name', None)
    if not expected_name:
        logger.error(f'No expected search engine name provided in rules: {expected_state}')
        return 0.0
    logger.info(f'Current search engine: {result_state}')
    logger.info(f'Expected search engine: {expected_name}')
    if result_state.strip().lower() == expected_name.strip().lower():
        logger.info('Search engine matches expected value')
        return 1.0
    else:
        logger.warning(f"Search engine mismatch: got '{result_state}', expected '{expected_name}'")
        return 0.0

def is_expected_url_pattern_match__907d006d1310a883a67bb7931a0dbae9(result, expected, **options) -> float:
    """
    Checks if the active tab URL matches multiple regex patterns.

    Args:
        result: The active tab info (string URL or dict with 'url' field)
        expected: Dictionary with 'expected' key containing list of regex patterns
        **options: Additional options

    Returns:
        float: 1.0 if all patterns match, 0.0 otherwise
    """
    if not result:
        logger.info('[PATTERN_MATCH] No result provided')
        return 0.0
    if isinstance(result, str):
        result_url = result
    elif isinstance(result, dict) and 'url' in result:
        result_url = result['url']
    else:
        logger.error(f'[PATTERN_MATCH] Invalid result format: {type(result)}')
        return 0.0
    logger.info(f'[PATTERN_MATCH] Result URL: {result_url}')
    patterns = expected.get('expected', [])
    logger.info(f'[PATTERN_MATCH] Expected patterns: {patterns}')
    for pattern in patterns:
        match = re.search(pattern, result_url, re.IGNORECASE)
        if not match:
            logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' not found in URL")
            return 0.0
        logger.info(f"[PATTERN_MATCH] Pattern '{pattern}' matched")
    logger.info('[PATTERN_MATCH] All patterns matched successfully')
    return 1.0

def check_webext_manifest__c1aa21b2(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_empty_description = expected.get('has_empty_description', False)
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_empty_description:
        checks += 1
        description = manifest.get('description')
        if description is None or description == '':
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest:
            bg = manifest['background']
            if 'scripts' in bg or 'service_worker' in bg or 'page' in bg:
                score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_extension_name__e5eabc9a(result, rules) -> float:
    """Check if the extension name matches the expected value.

    Args:
        result: Extension name from getter
        rules: Dict with 'expected' key containing expected extension name

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    expected_name = rules.get('expected')
    if result == expected_name:
        return 1.0
    return 0.0

def check_table_width__f4eb9543(result, expected, **options):
    """Check if table width meets minimum requirement.

    Args:
        result: Table position dict from getter
        expected: Dict with 'min_width' threshold
        **options: Additional options

    Returns:
        float: 1.0 if table width >= min_width, 0.0 otherwise
    """
    if result is None:
        return 0.0
    min_width = expected.get('min_width', 0)
    if result['width'] >= min_width:
        return 1.0
    else:
        return 0.0

def check_all_urls_contain__82d38938(result_state, expected_state, **options):
    """Check if all URLs contain a specific domain and meet minimum count requirements.

    This function verifies that:
    1. The URLs contain a specific domain (e.g., 'arxiv.org')
    2. There are at least a minimum number of valid URLs

    Args:
        result_state: List of URLs extracted from the document
        expected_state: Dict containing:
            - domain: Required domain string (e.g., 'arxiv.org')
            - min_count: Minimum number of URLs required
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
            - 1.0 if all URLs contain domain AND count >= min_count
            - Proportional score if count < min_count (count / min_count)
            - 0.0 if result_state is empty or invalid
    """
    if not result_state:
        return 0.0
    if not isinstance(result_state, list):
        return 0.0
    required_domain = expected_state.get('domain', 'arxiv.org')
    min_count = expected_state.get('min_count', 5)
    valid_urls = [url for url in result_state if required_domain in url]
    valid_count = len(valid_urls)
    if valid_count >= min_count:
        return 1.0
    elif valid_count > 0:
        return valid_count / min_count
    else:
        return 0.0

def check_chrome_setting__8cfb2b13(result, expected, **options):
    """
    Check if the Chrome Safe Browsing setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_tabs_and_bookmark__7a5a7856f1b642a4ade91ca81ca0f263000420251221151547(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the correct tab remains open and the URL is bookmarked.
    Awards partial credit: 0.5 for correct tab state, 0.5 for bookmark.

    Args:
        result: Dict with 'tabs' (list of tab dicts) and 'bookmarks' (bookmarks structure)
        expected: Dict with 'remaining_url' and 'bookmarked_url' keys

    Returns:
        float: Score between 0.0 and 1.0 (partial credit possible)
    """
    if not result:
        logger.info('No result data available')
        return 0.0
    tabs = result.get('tabs', [])
    bookmarks = result.get('bookmarks', {})
    remaining_url = expected.get('remaining_url')
    bookmarked_url = expected.get('bookmarked_url')
    if not remaining_url or not bookmarked_url:
        logger.error('Missing remaining_url or bookmarked_url in expected config')
        return 0.0
    score = 0.0
    logger.info(f'Checking tabs. Found {len(tabs)} open tabs')
    if len(tabs) == 1:
        tab_url = tabs[0].get('url', '')
        logger.info(f'Single tab URL: {tab_url}')
        if tab_url == remaining_url:
            logger.info(f'Correct tab remains open: {remaining_url}')
            score += 0.5
        else:
            logger.info(f'Wrong tab remains. Expected {remaining_url}, got {tab_url}')
    else:
        logger.info(f'Expected 1 tab, found {len(tabs)}')
    bookmark_bar = bookmarks.get('bookmark_bar', {})
    children = bookmark_bar.get('children', [])
    bookmark_urls = [child.get('url') for child in children if child.get('type') == 'url']
    logger.info(f'Bookmark bar URLs: {bookmark_urls}')
    if bookmarked_url in bookmark_urls:
        logger.info(f'URL is bookmarked: {bookmarked_url}')
        score += 0.5
    else:
        logger.info(f'URL not found in bookmarks: {bookmarked_url}')
    logger.info(f'Final score: {score}')
    return score

def check_chrome_ext_productivity__3ce2c199effd9eefc89344b99f4cd5a7(result: List[str], expected: Dict, **options) -> float:
    """Check if expected Chrome extensions are installed.

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' key containing list of required extension names
        **options: Additional options (not used)

    Returns:
        1.0 if all expected extensions are installed, 0.0 otherwise
    """
    if not result:
        logger.info('No extensions found in result')
        return 0.0
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        logger.warning('No expected extensions specified')
        return 0.0
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Expected extensions: {expected_extensions}')
    set_expected = set(expected_extensions)
    set_installed = set(result)
    if set_expected.issubset(set_installed):
        logger.info('All expected extensions are installed')
        return 1.0
    else:
        missing = set_expected - set_installed
        logger.info(f'Missing extensions: {missing}')
        return 0.0

def check_googledrive_file_count_and_names__dd26081a(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if Google Drive folder has the expected number of files with correct naming pattern

    Args:
        result: List of filenames from getter
        expected: Dict with:
            - count: Expected number of files
            - naming_pattern: Expected naming pattern (e.g., "N.png" where N is a number)
            - extension: Expected file extension (e.g., ".png")

    Returns:
        Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_count = expected.get('count', 0)
    naming_pattern = expected.get('naming_pattern', 'number')
    extension = expected.get('extension', '.png')
    actual_count = len(result)
    if actual_count != expected_count:
        logger.info(f'File count mismatch: expected {expected_count}, got {actual_count}')
        return 0.0
    if naming_pattern == 'number':
        expected_names = [f'{i}{extension}' for i in range(1, expected_count + 1)]
        if sorted(result) == sorted(expected_names):
            logger.info(f'All {expected_count} files have correct naming pattern')
            return 1.0
        else:
            logger.info(f'Naming pattern mismatch: expected {expected_names}, got {result}')
            return 0.0
    all_correct_ext = all((f.endswith(extension) for f in result))
    if all_correct_ext and len(result) == expected_count:
        return 1.0
    else:
        return 0.0

def is_extension_installed__d374210b(actual: str, rules: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {rules['type']}")

def check_webext_manifest__2ec2a9a8(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    expected_description_empty = expected.get('description_empty', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    if expected_description_empty:
        checks += 1
        description = manifest.get('description', '')
        if not description or description == '':
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_url__5525d1c8(result, expected, **options):
    """Check if URL is present in text.

    Args:
        result: Text from getter
        expected: Expected dict with 'url' key

    Returns:
        float: 1.0 if URL present, 0.0 otherwise
    """
    if not isinstance(result, str):
        return 0.0
    url = expected.get('url', '')
    if url and url in result:
        return 1.0
    return 0.0

def check_product_revenue_table__116fbd3460dada9be380bc664e224d99(result: Dict[str, float], expected: Dict[str, Any], **options) -> float:
    """
    Check if the product revenue table matches expected values.

    Args:
        result: Dict with product names as keys and revenue as values (from getter)
        expected: Dict with expected product revenue data (from evaluator.expected.rules)
        **options: Additional options (tolerance for numeric comparison)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_products = expected.get('products', {})
    tolerance = options.get('tolerance', 0.5)
    if not expected_products:
        return 0.0
    total_products = len(expected_products)
    if total_products == 0:
        return 0.0
    matched_products = 0
    for (product, expected_revenue) in expected_products.items():
        if product in result:
            actual_revenue = result[product]
            if abs(actual_revenue - expected_revenue) <= tolerance:
                matched_products += 1
    return matched_products / total_products

def check_gdrive_files__3dcb78c90ac64690f9f399090d59db08(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected files exist in Google Drive folder.

    Args:
        result: List of filenames found in Google Drive folder
        expected: Dict with 'expected_files' key containing list of required filenames
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_files = expected.get('expected_files', [])
    if not expected_files:
        logger.warning('No expected files specified')
        return 0.0
    if not result:
        logger.warning('No files found in Google Drive')
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
            logger.info(f'Found expected file: {expected_file}')
        else:
            logger.warning(f'Missing expected file: {expected_file}')
    score = found_count / len(expected_files)
    logger.info(f'Found {found_count}/{len(expected_files)} expected files. Score: {score}')
    return score

def is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_14_task_verify_1(result_state: Dict[str, Any], expected_state: Dict[str, Any], **options) -> float:
    """
    Checks if the expected bookmarks are in Chrome.
    This is a custom version that properly uses the 'names' field from the rule
    to determine the folder name instead of hardcoding 'Liked Authors'.

    Args:
        result_state: The bookmarks data from the getter
        expected_state: The expected state containing the rule
        **options: Additional options

    Returns:
        float: Score (1.0 if all bookmarks match, 0.0 otherwise)
    """
    bookmarks = result_state
    rule = expected_state
    if not bookmarks:
        return 0.0
    if rule['type'] == 'liked_authors_websites_urls':
        expected_folder_names = rule.get('names', [])
        if not expected_folder_names:
            logger.error('No folder names specified in rule')
            return 0.0
        expected_folder_name = expected_folder_names[0]
        if 'bookmark_bar' not in bookmarks or 'children' not in bookmarks['bookmark_bar']:
            logger.error('No bookmark_bar or children found in bookmarks')
            return 0.0
        target_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark.get('type') == 'folder' and bookmark.get('name') == expected_folder_name), None)
        if target_folder:
            logger.info(f"'{expected_folder_name}' folder exists")
            folder_urls = [bookmark['url'] for bookmark in target_folder.get('children', []) if bookmark.get('type') == 'url']
            logger.info(f"Here is the '{expected_folder_name}' folder's urls: {folder_urls}")
            urls = rule['urls']
            normalized_urls = []
            for url in urls:
                if isinstance(url, str):
                    normalized_urls.append([url])
                else:
                    normalized_urls.append(url)
            combinations = product(*normalized_urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    return 1.0
            return 0.0
        else:
            logger.info(f"'{expected_folder_name}' folder not found")
            return 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_chrome_tabs_and_bookmark__6d016e82(result: dict, expected: dict, **options):
    """Check Chrome tab count and bookmark.

    Args:
        result: Dict from getter with open_tab_count and bookmarks
        expected: Dict with expected values (from rules field)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        return 0.0
    expected_tab_count = expected.get('open_tab_count', 1)
    expected_has_bookmark = expected.get('has_stackoverflow_bookmark', True)
    expected_bookmark_title = expected.get('bookmark_title', 'Stack Overflow Reference')
    score = 0.0
    if result.get('open_tab_count') == expected_tab_count:
        score += 0.5
        logger.info(f"Tab count check passed: {result.get('open_tab_count')}")
    else:
        logger.info(f"Tab count check failed: got {result.get('open_tab_count')}, expected {expected_tab_count}")
    bookmarks = result.get('bookmarks', [])
    if expected_has_bookmark:
        found_bookmark = False
        for bookmark in bookmarks:
            if isinstance(bookmark, dict):
                url = bookmark.get('url', '')
                title = bookmark.get('name', '')
                if 'stackoverflow.com' in url and expected_bookmark_title in title:
                    found_bookmark = True
                    logger.info(f'Found matching bookmark: {title}')
                    break
        if found_bookmark:
            score += 0.5
        else:
            logger.info(f"Stackoverflow bookmark with title '{expected_bookmark_title}' not found")
    else:
        score += 0.5
    return score

def check_gdrive_files__fd7bb036(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_webext_manifest__c41914cdb60620ceca4e49be63e9e04c(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if web extension manifest.json matches expected structure.

    Args:
        result: Manifest JSON dict from getter (or None if file not found)
        expected: Expected manifest structure with rules to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.warning('Manifest file not found')
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_options_page = expected.get('has_options_page', False)
    score = 0.0
    total_checks = 0
    if expected_name is not None:
        total_checks += 1
        if result.get('name') == expected_name:
            score += 1
            logger.info(f'Name matches: {expected_name}')
        else:
            logger.warning(f"Name mismatch: expected '{expected_name}', got '{result.get('name')}'")
    if expected_version is not None:
        total_checks += 1
        if result.get('version') == expected_version:
            score += 1
            logger.info(f'Version matches: {expected_version}')
        else:
            logger.warning(f"Version mismatch: expected '{expected_version}', got '{result.get('version')}'")
    total_checks += 1
    has_background = 'background' in result and 'scripts' in result['background'] and (len(result['background']['scripts']) > 0)
    if expected_has_background == has_background:
        score += 1
        if has_background:
            logger.info(f"Background scripts found as expected: {result['background']['scripts']}")
        else:
            logger.info('Background scripts absent as expected')
    elif expected_has_background:
        logger.warning('Background scripts not found but expected')
    else:
        logger.warning('Background scripts found but should be absent')
    total_checks += 1
    has_browser_action = 'browser_action' in result
    if expected_has_browser_action == has_browser_action:
        score += 1
        if has_browser_action:
            logger.info('Browser action found as expected')
        else:
            logger.info('Browser action absent as expected')
    elif expected_has_browser_action:
        logger.warning('Browser action not found but expected')
    else:
        logger.warning('Browser action found but should be absent')
    total_checks += 1
    has_content_scripts = 'content_scripts' in result and len(result['content_scripts']) > 0
    if expected_has_content_scripts == has_content_scripts:
        score += 1
        if has_content_scripts:
            logger.info(f"Content scripts found as expected: {len(result['content_scripts'])} entries")
        else:
            logger.info('Content scripts absent as expected')
    elif expected_has_content_scripts:
        logger.warning('Content scripts not found but expected')
    else:
        logger.warning('Content scripts found but should be absent')
    total_checks += 1
    has_options_page = 'options_page' in result or 'options_ui' in result
    if expected_has_options_page == has_options_page:
        score += 1
        if has_options_page:
            logger.info('Options page found as expected')
        else:
            logger.info('Options page absent as expected')
    elif expected_has_options_page:
        logger.warning('Options page not found but expected')
    else:
        logger.warning('Options page found but should be absent')
    if total_checks == 0:
        logger.warning('No checks to perform')
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks} checks passed)')
    return final_score

def check_gdrive_pdf_in_folder__9b0aebe9(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_grf_table__d24a6b5d9bcf09db9cd5a9cf07143d27(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if GRF table has correct structure and year values.

    Args:
        result: Table data from getter with 'years' list and 'row_count'
        expected: Expected structure with 'min_rows' and 'years' list
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or not expected:
        return 0.0
    actual_years = result.get('years', [])
    actual_row_count = result.get('row_count', 0)
    expected_min_rows = expected.get('min_rows', 0)
    expected_years = expected.get('years', [])
    if actual_row_count < expected_min_rows:
        return 0.0
    actual_years_set = set(actual_years)
    expected_years_set = set(expected_years)
    if not expected_years_set:
        return 1.0 if actual_row_count >= expected_min_rows else 0.0
    matched_years = actual_years_set.intersection(expected_years_set)
    match_ratio = len(matched_years) / len(expected_years_set)
    if match_ratio >= 1.0:
        return 1.0
    elif match_ratio >= 0.8:
        return match_ratio
    else:
        return 0.0

def check_chrome_setting__1be3beb2(result, expected, **options):
    """
    Check if the Chrome Download location setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_extension_manifest__986ae666fddb67409d6f0937ecd8b146(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate browser extension manifest.json structure.

    Args:
        result: Manifest data from getter
        expected: Rules dict containing expected values for name, version, background, browser_action
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct field)
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dictionary')
        return 0.0
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    check_background = expected.get('check_background', False)
    check_browser_action = expected.get('check_browser_action', False)
    score = 0.0
    total_checks = 0
    total_checks += 1
    if result.get('name') == expected_name:
        score += 0.25
        logger.info(f'Name matches: {expected_name}')
    else:
        logger.info(f"Name mismatch: expected={expected_name}, got={result.get('name')}")
    total_checks += 1
    if result.get('version') == expected_version:
        score += 0.25
        logger.info(f'Version matches: {expected_version}')
    else:
        logger.info(f"Version mismatch: expected={expected_version}, got={result.get('version')}")
    if check_background:
        total_checks += 1
        background = result.get('background', {})
        if isinstance(background, dict) and 'scripts' in background:
            scripts = background.get('scripts', [])
            if isinstance(scripts, list) and len(scripts) > 0:
                score += 0.25
                logger.info(f'Background scripts found: {scripts}')
            else:
                logger.info('Background scripts field exists but is empty')
        else:
            logger.info('Background section missing or invalid')
    if check_browser_action:
        total_checks += 1
        browser_action = result.get('browser_action', {})
        if isinstance(browser_action, dict) and len(browser_action) > 0:
            has_fields = any((k in browser_action for k in ['default_popup', 'default_icon', 'default_title']))
            if has_fields:
                score += 0.25
                logger.info(f'Browser action found with fields: {list(browser_action.keys())}')
            else:
                logger.info('Browser action exists but missing expected fields')
        else:
            logger.info('Browser action section missing or invalid')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_url_match__3b2adddad1034fb6d584d5542b1c7034(result, expected, **options):
    """
    Check if the URL matches the expected value (case-insensitive, flexible matching).

    Args:
        result: str URL from getter
        expected: str expected URL

    Returns:
        float: 1.0 if match, 0.0 otherwise
    """
    if not isinstance(result, str) or not isinstance(expected, str):
        return 0.0
    result_normalized = result.lower().rstrip('/')
    expected_normalized = expected.lower().rstrip('/')
    return 1.0 if result_normalized == expected_normalized else 0.0

def check_googledrive_folder__a18b8359(result, expected, **options):
    """Check Google Drive folder information.

    Args:
        result: dict with "exists" and "file_count" from getter
        expected: dict with "exists" and "min_files"

    Returns:
        float: 1.0 if folder exists with min files, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    expected_exists = expected.get('exists', True)
    min_files = expected.get('min_files', 0)
    if result.get('exists', False) == expected_exists:
        if result.get('file_count', 0) >= min_files:
            return 1.0
    return 0.0

def check_webext_manifest__60d86cd5(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_channel_revenue_table__22652fc14c0135d2f835816b613f7f8c(result: Dict[str, float], expected: Dict[str, Any], **options) -> float:
    """
    Check if the sales channel revenue table matches expected values.

    Args:
        result: Dict with channel names as keys and revenue as values (from getter)
        expected: Dict with expected channel revenue data (from evaluator.expected.rules)
        **options: Additional options (tolerance for numeric comparison)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_channels = expected.get('channels', {})
    tolerance = options.get('tolerance', 0.5)
    if not expected_channels:
        return 0.0
    total_channels = len(expected_channels)
    if total_channels == 0:
        return 0.0
    matched_channels = 0
    for (channel, expected_revenue) in expected_channels.items():
        if channel in result:
            actual_revenue = result[channel]
            if abs(actual_revenue - expected_revenue) <= tolerance:
                matched_channels += 1
    return matched_channels / total_channels

def check_extension_enabled__1e533eef(result_state, expected_state, **options):
    """
    Check if a Chrome extension is both installed and enabled.

    This metric verifies that:
    1. The extension is installed
    2. The extension is enabled (state=1 in Chrome preferences)

    Args:
        result_state: Dictionary from getter containing:
            - extension_name (str): The name of the extension
            - is_installed (bool): Whether the extension is installed
            - is_enabled (bool): Whether the extension is enabled
        expected_state: Dictionary containing:
            - extension_name (str): Expected extension name
            - should_be_enabled (bool): Whether extension should be enabled

    Returns:
        float: 1.0 if extension is installed and enabled as expected, 0.0 otherwise
    """
    if not result_state:
        logger.error('[CHECK_EXTENSION_ENABLED] No result state provided')
        return 0.0
    if not expected_state:
        logger.error('[CHECK_EXTENSION_ENABLED] No expected state provided')
        return 0.0
    extension_name = result_state.get('extension_name', '')
    is_installed = result_state.get('is_installed', False)
    is_enabled = result_state.get('is_enabled', False)
    expected_name = expected_state.get('extension_name', '')
    should_be_enabled = expected_state.get('should_be_enabled', True)
    logger.info(f'[CHECK_EXTENSION_ENABLED] Checking extension: {extension_name}')
    logger.info(f'[CHECK_EXTENSION_ENABLED] Expected name: {expected_name}')
    logger.info(f'[CHECK_EXTENSION_ENABLED] Is installed: {is_installed}')
    logger.info(f'[CHECK_EXTENSION_ENABLED] Is enabled: {is_enabled}')
    logger.info(f'[CHECK_EXTENSION_ENABLED] Should be enabled: {should_be_enabled}')
    if extension_name != expected_name:
        logger.error(f"[CHECK_EXTENSION_ENABLED] Extension name mismatch: '{extension_name}' != '{expected_name}'")
        return 0.0
    if not is_installed:
        logger.error(f"[CHECK_EXTENSION_ENABLED] Extension '{extension_name}' is not installed")
        return 0.0
    if is_enabled != should_be_enabled:
        logger.error(f'[CHECK_EXTENSION_ENABLED] Extension enabled state mismatch: {is_enabled} != {should_be_enabled}')
        return 0.0
    logger.info(f"[CHECK_EXTENSION_ENABLED] Extension '{extension_name}' is properly installed and enabled")
    return 1.0

def check_file_executable__2bea57f7(result, expected, **options):
    """Check if file has executable permissions.

    Args:
        result: File permissions string
        expected: Dict (not used, checking for 'x' permission)
        **options: Additional options

    Returns:
        float: 1.0 if file is executable by user, 0.0 otherwise
    """
    if result is None:
        return 0.0
    if len(result) >= 4 and result[3] == 'x':
        return 1.0
    else:
        return 0.0

def check_chrome_setting__5936c775(result, expected, **options):
    """
    Check if the Chrome Minimum font size setting is within expected range.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type', 'min', and 'max' specification

    Returns:
        float: 1.0 if in range, 0.0 otherwise
    """
    setting_value = result.get('setting_value', 0)
    if expected['type'] == 'range':
        min_val = expected.get('min', float('-inf'))
        max_val = expected.get('max', float('inf'))
        return 1.0 if min_val <= setting_value <= max_val else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_csv_table_phys__5e1848a9(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_recreation_url__6dc3893f96ccd943c500af1756962de6(result: Dict[str, str], expected: Dict[str, Any], **options) -> float:
    """
    Check if the current URL matches the expected recreation.gov pattern.

    Args:
        result: Dictionary with 'url' key
        expected: Expected configuration with url_pattern field
        **options: Additional options

    Returns:
        float: 1.0 if URL matches pattern, 0.0 otherwise
    """
    logger.info(f'[CHECK_RECREATION_URL] Result: {result}')
    logger.info(f'[CHECK_RECREATION_URL] Expected: {expected}')
    if not result or 'url' not in result:
        logger.info('[CHECK_RECREATION_URL] Result is empty or missing URL, returning 0.0')
        return 0.0
    actual_url = result['url']
    url_pattern = expected.get('url_pattern', '')
    if re.search(url_pattern, actual_url, re.IGNORECASE):
        logger.info(f"[CHECK_RECREATION_URL] URL matches pattern: '{url_pattern}' matches '{actual_url}'")
        return 1.0
    else:
        logger.info(f"[CHECK_RECREATION_URL] URL mismatch: pattern '{url_pattern}' does not match '{actual_url}'")
        return 0.0

def check_gdrive_pdf_in_folder__a29f27f0(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_prettier_extension__4d8ac023(actual: str, rules: Dict, **options):
    """
    Check if Prettier extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_gdrive_pdf__75da4280(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_chrome_setting__fd2ee811(result, expected, **options):
    """
    Check if the Chrome Autofill setting matches the expected value.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type' and 'value' specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    setting_value = result.get('setting_value')
    if expected['type'] == 'value':
        return 1.0 if setting_value == expected['value'] else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_webext_manifest__e32cecfb(result, expected, **options):
    """
    Check if manifest.json has correct structure and values.

    Args:
        result: Path to manifest.json file
        expected: Dictionary with expected manifest fields
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        with open(result, 'r') as f:
            manifest = json.load(f)
    except (json.JSONDecodeError, IOError):
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_browser_action = expected.get('has_browser_action', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_has_options = expected.get('has_options', False)
    expected_has_blank_description = expected.get('has_blank_description', False)
    score = 0.0
    checks = 0
    if expected_name is not None:
        checks += 1
        if manifest.get('name') == expected_name:
            score += 1.0
    if expected_version is not None:
        checks += 1
        if manifest.get('version') == expected_version:
            score += 1.0
    if expected_has_background:
        checks += 1
        if 'background' in manifest and 'scripts' in manifest['background']:
            score += 1.0
    if expected_has_browser_action:
        checks += 1
        if 'browser_action' in manifest:
            score += 1.0
    if expected_has_content_scripts:
        checks += 1
        if 'content_scripts' in manifest:
            score += 1.0
    if expected_has_page_action:
        checks += 1
        if 'page_action' in manifest:
            score += 1.0
    if expected_has_options:
        checks += 1
        if 'options_ui' in manifest or 'options_page' in manifest:
            score += 1.0
    if expected_has_blank_description:
        checks += 1
        description = manifest.get('description', '')
        if not description or description.strip() == '':
            score += 1.0
    return score / checks if checks > 0 else 0.0

def check_gdrive_files__839d38d5(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_extension_version__f107c64c583fbfb012ea91827c8f61e3(result, expected, **options):
    """
    Check if the extension version matches the expected value.

    Args:
        result: String from getter containing the extension version
        expected: Dict with 'extension_version' key
        **options: Additional options (unused)

    Returns:
        float: 1.0 if versions match, 0.0 otherwise
    """
    expected_version = expected.get('extension_version', '')
    if result == expected_version:
        return 1.0
    return 0.0

def check_extension_version__ae6416e4(result: Dict[str, str], expected: Dict[str, Any], **options) -> float:
    """Check if a specific extension has the expected version.

    Args:
        result: Dict mapping extension names to version strings
        expected: Rules dict with 'extension_name' and 'version' keys
        **options: Additional options

    Returns:
        1.0 if extension has the expected version, 0.0 otherwise
    """
    extension_name = expected.get('extension_name', '')
    expected_version = expected.get('version', '')
    if not extension_name or not expected_version:
        return 0.0
    if extension_name in result and result[extension_name] == expected_version:
        return 1.0
    else:
        return 0.0

def check_extension_details__ae6416e4(result: List[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if an extension with specific properties exists.

    Args:
        result: List of extension details
        expected: Rules dict with extension properties to match
        **options: Additional options

    Returns:
        1.0 if matching extension found, 0.0 otherwise
    """
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    if not expected_name:
        return 0.0
    for ext in result:
        if ext.get('name') == expected_name:
            if expected_version:
                if ext.get('version') == expected_version:
                    return 1.0
            else:
                return 1.0
    return 0.0

def check_extension_manifest_version__b7035c23581ef82d8725cee1b3aa987f(result, expected, **options):
    """
    Check if the extension manifest version matches the expected value.

    Args:
        result: Integer from getter containing the manifest_version
        expected: Dict with 'manifest_version' key
        **options: Additional options (unused)

    Returns:
        float: 1.0 if manifest versions match, 0.0 otherwise
    """
    expected_version = expected.get('manifest_version', 0)
    if result == expected_version:
        return 1.0
    return 0.0

def check_chrome_ext_last_two__f117f7ab324ab80f0f8ad254a42cb210(result: List[str], expected: Dict, baseline: List[str]=None, **options) -> float:
    """Check if expected Chrome extensions were newly installed during task execution.

    Args:
        result: List of installed extension names after task execution
        expected: Dict with 'expected' key containing list of required extension names
        baseline: List of installed extension names before task execution (optional)
        **options: Additional options (not used)

    Returns:
        1.0 if the expected extensions were newly added, 0.0 otherwise
    """
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        logger.warning('No expected extensions specified')
        return 0.0
    logger.info(f'Baseline extensions: {(baseline if baseline else [])}')
    logger.info(f'Result extensions: {(result if result else [])}')
    logger.info(f'Expected extensions: {expected_extensions}')
    if baseline is None:
        logger.warning('No baseline provided, using simple subset check')
        if not result:
            logger.info('No extensions found in result')
            return 0.0
        set_expected = set(expected_extensions)
        set_installed = set(result)
        if set_expected.issubset(set_installed):
            logger.info('All expected extensions are installed')
            return 1.0
        else:
            missing = set_expected - set_installed
            logger.info(f'Missing extensions: {missing}')
            return 0.0
    baseline_set = set(baseline) if baseline else set()
    result_set = set(result) if result else set()
    expected_set = set(expected_extensions)
    newly_added = result_set - baseline_set
    logger.info(f'Newly added extensions: {newly_added}')
    if expected_set.issubset(newly_added):
        logger.info('All expected extensions were newly installed')
        return 1.0
    else:
        pre_existing = expected_set & baseline_set
        missing = expected_set - result_set
        not_newly_added = expected_set - newly_added
        if pre_existing:
            logger.info(f'Extensions already installed before task (false positive): {pre_existing}')
        if missing:
            logger.info(f'Expected extensions not installed: {missing}')
        if not_newly_added and (not pre_existing):
            logger.info(f'Expected extensions not newly added: {not_newly_added}')
        return 0.0

def check_html_exists__7b7b0b2f(result_state: Optional[Dict[str, Any]], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if an HTML file exists as expected and contains valid HTML content.

    Args:
        result_state: Dict from getter with 'exists', 'is_valid_html', 'is_non_empty' keys
                     (None if there was an error)
        expected_state: Dict containing 'should_exist' key (when type='rule', this IS the rules dict)
        **options: Additional options (not used)

    Returns:
        float: 1.0 if the file exists, is valid HTML, and is non-empty (when expected), 0.0 otherwise
    """
    if result_state is None:
        return 0.0
    should_exist = expected_state.get('should_exist', True)
    if not should_exist:
        return 1.0 if not result_state.get('exists', False) else 0.0
    if result_state.get('exists', False) and result_state.get('is_valid_html', False) and result_state.get('is_non_empty', False):
        return 1.0
    else:
        return 0.0

def check_gdrive_pattern__6538c56492a239416edab32324471fe6(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if Google Drive files meet naming pattern criteria with exact count and sequential numbering.

    Args:
        result: Dict from getter with keys: matching_count, matching_files, file_numbers, all_files
        expected: Dict with validation rules:
            - exact_count: int - exact number of files that must match pattern
            - verify_sequential: bool - whether to verify sequential numbering from 1 to N

    Returns:
        float: 1.0 if all criteria met, 0.0 otherwise
    """
    if not isinstance(result, dict):
        return 0.0
    matching_count = result.get('matching_count', 0)
    file_numbers = result.get('file_numbers', [])
    exact_count = expected.get('exact_count')
    verify_sequential = expected.get('verify_sequential', False)
    if exact_count is not None:
        if matching_count != exact_count:
            return 0.0
    if verify_sequential and exact_count is not None:
        expected_sequence = list(range(1, exact_count + 1))
        if file_numbers != expected_sequence:
            return 0.0
    return 1.0

def check_recreation_search__e0d56e1bf714f2c8f287d0502b986b63(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """
    Check if the recreation.gov page has a search box.

    Args:
        result: Dictionary with 'has_search_box' and 'count' keys
        expected: Expected configuration with has_search field
        **options: Additional options

    Returns:
        float: 1.0 if search box presence matches expectation, 0.0 otherwise
    """
    logger.info(f'[CHECK_RECREATION_SEARCH] Result: {result}')
    logger.info(f'[CHECK_RECREATION_SEARCH] Expected: {expected}')
    if not result:
        logger.info('[CHECK_RECREATION_SEARCH] Result is empty, returning 0.0')
        return 0.0
    actual_has_search = result.get('has_search_box', False)
    expected_has_search = expected.get('has_search', True)
    if actual_has_search == expected_has_search:
        logger.info(f'[CHECK_RECREATION_SEARCH] Search box presence matches expectation: {actual_has_search}')
        return 1.0
    else:
        logger.info(f'[CHECK_RECREATION_SEARCH] Search box presence mismatch: expected {expected_has_search}, got {actual_has_search}')
        return 0.0

def check_table_bottom_left__6a4e1dd4(result, expected, **options):
    """Check if table is positioned in bottom-left corner.

    Args:
        result: Table position dict from getter
        expected: Dict with 'min_top' and 'max_left'
        **options: Additional options

    Returns:
        float: 1.0 if both conditions met, 0.0 otherwise
    """
    if result is None:
        return 0.0
    min_top = expected.get('min_top', 0)
    max_left = expected.get('max_left', 0)
    if result['top'] > min_top and result['left'] < max_left:
        return 1.0
    else:
        return 0.0

def check_gdrive_pdf__c340b25e(result: str, expected: Dict[str, Any], **options) -> float:
    """Check if PDF file meets expected criteria.
    
    Args:
        result: Path to PDF file (from getter)
        expected: Expected values dict from expected.rules with:
            - relation: Comparison operator (e.g., 'eq', 'ge', 'le')
            - ref_value: Expected page count
            - verify_content: Whether to verify PDF has actual content
            
    Returns:
        float: 1.0 if PDF meets all criteria, 0.0 otherwise
    """
    import operator
    if result is None or not os.path.exists(result):
        return 0.0
    try:
        reader = PdfReader(result)
        nb_pages = len(reader.pages)
        relation = expected.get('relation', 'eq')
        ref_value = expected.get('ref_value', 1)
        page_count_matches = getattr(operator, relation)(nb_pages, ref_value)
        if not page_count_matches:
            return 0.0
        if expected.get('verify_content', False):
            file_size = os.path.getsize(result)
            if file_size < 10240:
                return 0.0
        return 1.0
    except Exception as e:
        print(f'Error checking PDF: {e}')
        return 0.0

def check_cpp_extension__d45d3417(actual: str, rules: Dict, **options):
    """
    Check if C/C++ extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        rules: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if rules['type'] == 'contain':
        if rules['expected'] in actual:
            return 1.0
        return 0.0
    elif rules['type'] == 'not_contain':
        if rules['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {rules['type']}")

def check_git_remote_url__f8582c17(result, expected, **options):
    """
    Check if git remote URL contains expected substring.

    Args:
        result: Remote URL from getter
        expected: Expected URL substring

    Returns:
        float: 1.0 if URL contains expected substring, 0.0 otherwise
    """
    expected_url = expected.get('url', '')
    if expected_url in result:
        return 1.0
    return 0.0

def check_amoxicillin_url__b070486d(result, rules) -> float:
    """
    Check if the URL matches one of the expected Amoxicillin side effects URL patterns.
    This function implements OR logic - returns 1.0 if ANY of the patterns match.

    Args:
        result: The active URL (string or dict with 'url' field)
        rules: Dict containing 'expected' list with two URL patterns

    Returns:
        float: 1.0 if any pattern matches, 0.0 otherwise
    """
    if not result:
        return 0.0
    if isinstance(result, str):
        result_url = result
    elif isinstance(result, dict) and 'url' in result:
        result_url = result['url']
    else:
        logger.error(f"Invalid result format: {type(result)}, expected string URL or dict with 'url' field")
        return 0.0
    logger.info(f'Result URL to match: {result_url}')
    patterns = rules.get('expected', [])
    if not patterns:
        logger.error('No expected patterns provided in rules')
        return 0.0
    logger.info(f'Expected patterns: {patterns}')
    for pattern in patterns:
        try:
            match = re.search(pattern, result_url)
            if match:
                logger.info(f'Pattern matched: {pattern}')
                return 1.0
        except re.error as e:
            logger.error(f"Invalid regex pattern '{pattern}': {e}")
            continue
    logger.info('No patterns matched')
    return 0.0

def check_jupyter_extension__bb12eb1b(actual: str, expected: Dict, **options):
    """
    Check if Jupyter extension is installed in VSCode.

    Args:
        actual: Output from 'code --list-extensions' command
        expected: Dictionary with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if expected['type'] == 'contain':
        if expected['expected'] in actual:
            return 1.0
        return 0.0
    elif expected['type'] == 'not_contain':
        if expected['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown check type: {expected['type']}")

def is_expected_bookmarks__a82b78bb(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome with support for equal_contrib_authors_urls and liked_authors_websites_urls.

    This metric handles multiple rule types:
    - 'equal_contrib_authors_urls': Checks if a folder with name from rule['names'][0] contains URLs from rule['urls']
    - 'liked_authors_websites_urls': Checks if a folder with name from rule['names'][0] contains URLs from rule['urls']

    Args:
        bookmarks: Chrome bookmarks data structure
        rule: Rule dictionary with type, names, and urls

    Returns:
        1.0 if the expected bookmarks are found, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] in ['equal_contrib_authors_urls', 'liked_authors_websites_urls']:
        if not rule.get('names') or len(rule['names']) == 0:
            logger.error("Rule 'names' field is missing or empty")
            return 0.0
        expected_folder_name = rule['names'][0]
        target_folder = None
        for bookmark in bookmarks['bookmark_bar']['children']:
            if bookmark['type'] == 'folder' and bookmark['name'] == expected_folder_name:
                target_folder = bookmark
                break
        if not target_folder:
            logger.info(f"Folder '{expected_folder_name}' not found in bookmark bar")
            return 0.0
        logger.info(f"Folder '{expected_folder_name}' exists")
        folder_urls = [bookmark['url'] for bookmark in target_folder['children'] if bookmark['type'] == 'url']
        logger.info(f"URLs in '{expected_folder_name}' folder: {folder_urls}")
        urls = rule['urls']
        normalized_urls = []
        for url in urls:
            if isinstance(url, str):
                normalized_urls.append([url])
            else:
                normalized_urls.append(url)
        combinations = product(*normalized_urls)
        for combination in combinations:
            if set(combination) == set(folder_urls):
                logger.info(f'Found matching combination: {combination}')
                return 1.0
        logger.info(f'No matching URL combination found. Expected one of the combinations from {normalized_urls}, got {folder_urls}')
        return 0.0
    else:
        raise TypeError(f"{rule['type']} not supported by this metric!")

def check_gdrive_nested_files__688ade84b6705b21f2a27dc4863ba216(result: List[str], expected: Dict[str, Any], **options) -> float:
    """Check if expected files exist in nested Google Drive folder structure.

    Args:
        result: List of file paths that exist (formatted as "folder/subfolder/file.ext")
        expected: Dict with 'expected_paths' key containing list of required file paths
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    expected_paths = expected.get('expected_paths', [])
    if not expected_paths:
        logger.warning('No expected file paths specified')
        return 0.0
    if not result:
        logger.warning('No files found in Google Drive')
        return 0.0
    found_count = 0
    for expected_path in expected_paths:
        if expected_path in result:
            found_count += 1
            logger.info(f'Found expected file at: {expected_path}')
        else:
            logger.warning(f'Missing expected file at: {expected_path}')
    score = found_count / len(expected_paths)
    logger.info(f'Found {found_count}/{len(expected_paths)} expected file paths. Score: {score}')
    return score

def check_table_h_centered__dd5268aa(result, expected, **options):
    """Check if table is horizontally centered.

    Args:
        result: Table position dict from getter
        expected: Dict with 'target_left' and 'tolerance'
        **options: Additional options

    Returns:
        float: 1.0 if within tolerance, 0.0 otherwise
    """
    if result is None:
        return 0.0
    target_left = expected.get('target_left', 0)
    tolerance = expected.get('tolerance', 0)
    if abs(result['left'] - target_left) <= tolerance:
        return 1.0
    else:
        return 0.0

def check_extension__e13be972(result, expected, **options):
    """Check if file has expected extension.

    Args:
        result: Actual extension (str)
        expected: Dict with 'extension' key
        **options: Additional options

    Returns:
        float: 1.0 if extension matches, 0.0 otherwise
    """
    expected_ext = expected.get('extension', 'xlsx')
    if result.lower() == expected_ext.lower():
        return 1.0
    else:
        return 0.0

def check_month_revenue_table__321ed1a008442c1379822ece597d5a12(result: Dict[str, float], expected: Dict[str, Any], **options) -> float:
    """
    Check if the month revenue table matches expected values.

    Args:
        result: Dict with month names as keys and revenue as values (from getter)
        expected: Dict with expected month revenue data (from evaluator.expected.rules)
        **options: Additional options (tolerance for numeric comparison)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    expected_months = expected.get('months', {})
    tolerance = options.get('tolerance', 0.5)
    if not expected_months:
        return 0.0
    total_months = len(expected_months)
    if total_months == 0:
        return 0.0
    matched_months = 0
    for (month, expected_revenue) in expected_months.items():
        if month in result:
            actual_revenue = result[month]
            if abs(actual_revenue - expected_revenue) <= tolerance:
                matched_months += 1
    return matched_months / total_months

def check_table_row_count_by_index__83874f16(result, expected, **options):
    """Check table row count."""
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    result_index = result.get('table_index', -1)
    expected_index = expected.get('table_index', -1)
    if result_index != expected_index:
        return 0.0
    result_rows = result.get('row_count', 0)
    min_rows = expected.get('min_rows', 0)
    if result_rows >= min_rows:
        return 1.0
    else:
        return 0.0

def check_csv_table_med__1bf7a84e(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_extension_manifest_version__ae6416e4(result: Dict[str, Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """Check if a specific extension has the expected manifest version and is unpacked.

    Args:
        result: Dict mapping extension names to dicts with 'manifest_version' and 'is_unpacked' keys
        expected: Rules dict with 'extension_name' and 'manifest_version' keys
        **options: Additional options

    Returns:
        1.0 if extension has the expected manifest version and is unpacked, 0.0 otherwise
    """
    extension_name = expected.get('extension_name', '')
    expected_manifest_version = expected.get('manifest_version', 0)
    if not extension_name or not expected_manifest_version:
        return 0.0
    if extension_name in result:
        ext_info = result[extension_name]
        has_correct_version = ext_info.get('manifest_version') == expected_manifest_version
        is_unpacked = ext_info.get('is_unpacked', False)
        if has_correct_version and is_unpacked:
            return 1.0
    return 0.0

def check_chrome_setting__fca61186(result, expected, **options):
    """
    Check if the Chrome Fixed-width font size setting is within expected range.

    Args:
        result: Dictionary with 'setting_value' key (from getter)
        expected: Dictionary with 'type', 'min', and 'max' specification

    Returns:
        float: 1.0 if in range, 0.0 otherwise
    """
    setting_value = result.get('setting_value', 0)
    if expected['type'] == 'range':
        min_val = expected.get('min', float('-inf'))
        max_val = expected.get('max', float('inf'))
        return 1.0 if min_val <= setting_value <= max_val else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_html_conversion__b5ca523a54cc6ae837c08b60601febd6(result, expected, **options):
    """
    Check if all Word documents were correctly converted to .html format.

    Args:
        result: Dict with 'command_output', 'html_count', 'archive_exists' from getter
        expected: Dictionary of rules from config
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not isinstance(expected, dict):
        return 0.0
    command_output = result.get('command_output', '')
    html_count = result.get('html_count', 0)
    archive_exists = result.get('archive_exists', False)
    expected_command = expected.get('command_match', 'catch the desired command')
    expected_html_count = 17
    score = 0.0
    if isinstance(command_output, str) and expected_command in command_output:
        score += 0.5
    if archive_exists and html_count == expected_html_count:
        score += 0.5
    elif archive_exists and html_count > 0:
        score += 0.5 * (html_count / expected_html_count)
    return score

def is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_task_verify_1(bookmarks: Dict, rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in a specific folder in Chrome's bookmark bar.
    This version properly uses the 'names' field from the rule to check for the correct folder name.

    Args:
        bookmarks: Chrome bookmarks structure
        rule: Rule configuration with 'type', 'names' (folder names), and 'urls' (expected URLs)

    Returns:
        float: 1.0 if the folder exists with the correct URLs, 0.0 otherwise
    """
    if not bookmarks:
        return 0.0
    if rule['type'] == 'bookmark_bar_folders_with_urls':
        folder_names = rule.get('names', [])
        if not folder_names:
            logger.error('No folder names specified in rule')
            return 0.0
        folder_name = folder_names[0]
        target_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder' and bookmark['name'] == folder_name), None)
        if target_folder:
            logger.info(f"'{folder_name}' folder exists")
            folder_urls = [bookmark['url'] for bookmark in target_folder['children'] if bookmark['type'] == 'url']
            logger.info(f"Here is the '{folder_name}' folder's urls: {folder_urls}")
            urls = rule['urls']
            for (idx, url) in enumerate(urls):
                if isinstance(url, str):
                    urls[idx] = [url]
            combinations = product(*urls)
            for combination in combinations:
                if set(combination) == set(folder_urls):
                    return 1.0
            return 0.0
        else:
            logger.info(f"Folder '{folder_name}' not found in bookmark bar")
            return 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_extension_count__4329ad6c(result, expected, **options):
    """
    Check if at least min_count installed extensions are from the recommended list.

    Args:
        result: Dict with two keys:
            - 'installed': List of installed extension names
            - 'recommended': List of recommended extension names from document
        expected: Dict with 'rules' containing 'min_count' for minimum expected extensions from the list
        **options: Additional options

    Returns:
        float: 1.0 if at least min_count extensions from recommended list are installed, 0.0 otherwise
    """
    if isinstance(result, list):
        logger.warning('Received list format instead of dict - cannot verify against recommended list')
        installed_extensions = result
        recommended_extensions = []
    elif isinstance(result, dict):
        installed_extensions = result.get('installed', [])
        recommended_extensions = result.get('recommended', [])
    else:
        logger.error(f'Unexpected result type: {type(result)}')
        return 0.0
    if not installed_extensions:
        installed_extensions = []
    if not recommended_extensions:
        recommended_extensions = []
    if isinstance(expected, dict):
        min_count = expected.get('min_count', 0)
    else:
        min_count = 0
    logger.info(f'Installed extensions: {installed_extensions}')
    logger.info(f'Recommended extensions: {recommended_extensions}')
    matching_extensions = []
    for installed in installed_extensions:
        for recommended in recommended_extensions:
            if installed.lower().strip() == recommended.lower().strip():
                matching_extensions.append(installed)
                break
    matching_count = len(matching_extensions)
    logger.info(f'Matching extensions from recommended list: {matching_extensions}')
    logger.info(f'Matching count: {matching_count}, minimum required: {min_count}')
    if matching_count >= min_count:
        return 1.0
    else:
        return 0.0

def check_chrome_setting__a785d845(result, expected, **options):
    """
    Check if the Chrome zoom level setting matches the expected value or range.

    Args:
        result: Dictionary with 'zoom_level' key (from getter)
        expected: Dictionary with 'type' and value/range specification

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    zoom_level = result.get('zoom_level', 0.0)
    if expected['type'] == 'value':
        return 1.0 if zoom_level == expected['value'] else 0.0
    elif expected['type'] == 'range':
        min_val = expected.get('min', float('-inf'))
        max_val = expected.get('max', float('inf'))
        return 1.0 if min_val <= zoom_level <= max_val else 0.0
    else:
        logger.error(f"Unknown type: {expected['type']}")
        return 0.0

def check_csv_table_ee__c4163e13(result, expected, **options):
    """Check if CSV table matches expected structure and values.

    Args:
        result: List of rows from getter
        expected: Dict with 'headers' and 'data' keys
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not result or len(result) < 2:
        return 0.0
    expected_headers = expected.get('headers', [])
    expected_data = expected.get('data', [])
    score = 0.0
    if expected_headers:
        actual_headers = result[0] if len(result) > 0 else []
        if actual_headers == expected_headers:
            score += 0.2
    if expected_data:
        for (i, expected_row) in enumerate(expected_data):
            actual_row_idx = i + 1
            if actual_row_idx < len(result):
                actual_row = result[actual_row_idx]
                row_matches = True
                for (j, expected_cell) in enumerate(expected_row):
                    if j < len(actual_row):
                        actual_cell = actual_row[j].strip()
                        expected_cell_str = str(expected_cell).strip()
                        if '%' in expected_cell_str:
                            try:
                                actual_val = float(actual_cell.rstrip('%'))
                                expected_val = float(expected_cell_str.rstrip('%'))
                                if abs(actual_val - expected_val) > 0.01:
                                    row_matches = False
                                    break
                            except:
                                if actual_cell != expected_cell_str:
                                    row_matches = False
                                    break
                        elif actual_cell != expected_cell_str:
                            row_matches = False
                            break
                    else:
                        row_matches = False
                        break
                if row_matches:
                    score += 0.8 / len(expected_data)
    return min(score, 1.0)

def check_websites_filled__53b95efcf5f95ba5e65301e716b45c01(result, expected, **options):
    """Check if specific restaurants have valid website URLs filled.

    Args:
        result: Dict mapping restaurant names to website URLs
        expected: Dict with 'required_restaurants' key listing required restaurant names
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    required_restaurants = expected.get('required_restaurants', [])
    if not required_restaurants:
        return 0.0
    filled_count = 0
    for restaurant in required_restaurants:
        website = result.get(restaurant, '')
        if website and len(website) > 0:
            website_lower = website.lower()
            if any((pattern in website_lower for pattern in ['http', 'www.', '.com', '.net', '.org', '.hk'])):
                filled_count += 1
            elif '.' in website:
                filled_count += 1
    total_required = len(required_restaurants)
    if filled_count >= total_required:
        return 1.0
    else:
        return round(filled_count / total_required, 2)

def check_extension_manifest__842d772f81b13ef554c1bda7e1c59bb7(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate browser extension manifest.json structure.

    Args:
        result: Manifest data from getter
        expected: Rules dict containing expected values for name, version, background, browser_action
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0 (partial credit for each correct field)
    """
    if not result or not isinstance(result, dict):
        logger.warning('Result is empty or not a dictionary')
        return 0.0
    expected_name = expected.get('name', '')
    expected_version = expected.get('version', '')
    check_background = expected.get('check_background', False)
    check_browser_action = expected.get('check_browser_action', False)
    score = 0.0
    total_checks = 0
    total_checks += 1
    if result.get('name') == expected_name:
        score += 0.25
        logger.info(f'Name matches: {expected_name}')
    else:
        logger.info(f"Name mismatch: expected={expected_name}, got={result.get('name')}")
    total_checks += 1
    if result.get('version') == expected_version:
        score += 0.25
        logger.info(f'Version matches: {expected_version}')
    else:
        logger.info(f"Version mismatch: expected={expected_version}, got={result.get('version')}")
    if check_background:
        total_checks += 1
        background = result.get('background', {})
        if isinstance(background, dict) and 'scripts' in background:
            scripts = background.get('scripts', [])
            if isinstance(scripts, list) and len(scripts) > 0:
                score += 0.25
                logger.info(f'Background scripts found: {scripts}')
            else:
                logger.info('Background scripts field exists but is empty')
        else:
            logger.info('Background section missing or invalid')
    if check_browser_action:
        total_checks += 1
        browser_action = result.get('browser_action', {})
        if isinstance(browser_action, dict) and len(browser_action) > 0:
            has_fields = any((k in browser_action for k in ['default_popup', 'default_icon', 'default_title']))
            if has_fields:
                score += 0.25
                logger.info(f'Browser action found with fields: {list(browser_action.keys())}')
            else:
                logger.info('Browser action exists but missing expected fields')
        else:
            logger.info('Browser action section missing or invalid')
    logger.info(f'Final score: {score:.2f}')
    return score

def check_chrome_extension_manifest__3d480472b35ce7003ca943ba6b2307fa(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate Chrome extension manifest and files.

    Args:
        result: Dict from getter with 'manifest', 'files_exist', 'all_files_exist'
        expected: Dict with expected manifest structure rules:
            - manifest_checks: List of dicts with 'key' (list of nested keys) and 'value' (expected value)
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    manifest = result.get('manifest', {})
    all_files_exist = result.get('all_files_exist', False)
    file_score = 0.5 if all_files_exist else 0.0
    manifest_checks = expected.get('manifest_checks', [])
    if not manifest_checks:
        return file_score
    passed_checks = 0
    total_checks = len(manifest_checks)
    for check in manifest_checks:
        key_path = check.get('key', [])
        expected_value = check.get('value')
        current = manifest
        try:
            for key in key_path:
                current = current[key]
            if current == expected_value:
                passed_checks += 1
        except (KeyError, TypeError):
            pass
    manifest_score = 0.5 * (passed_checks / total_checks) if total_checks > 0 else 0.0
    return file_score + manifest_score

def check_gdrive_pdf_in_folder__adf6c9b9(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def check_chrome_third_party_cookies__1394774d(result, expected, **options):
    """Check if Chrome third-party cookies blocking is in expected state.

    Args:
        result: Third-party cookies settings dict from getter
        expected: Expected rules dict with 'block_third_party' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_block = expected.get('block_third_party', False)
    actual_block = result.get('block_third_party', False)
    logger.info(f'Expected block third-party cookies: {expected_block}')
    logger.info(f'Actual block third-party cookies: {actual_block}')
    return 1.0 if actual_block == expected_block else 0.0

def check_gdrive_folder__55cd0f01(result, expected, **options):
    """Check if folder exists with correct number of PDFs.

    Args:
        result: Folder info from getter
        expected: Expected folder configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict):
        return 0.0
    if not result.get('folder_exists', False):
        return 0.0
    expected_count = expected.get('file_count', 0)
    actual_count = result.get('file_count', 0)
    if actual_count == expected_count:
        return 1.0
    else:
        return min(actual_count / expected_count, 1.0) if expected_count > 0 else 0.0

def check_vim_tabstop__e2e649c5246b836f874ad28b723333bc(result: str, expected: Dict[str, List[str]]) -> float:
    """
    Check if the result contains expected include strings and excludes unwanted strings.

    Args:
        result: Output from the getter function
        expected: Dictionary with 'include' and 'exclude' lists

    Returns:
        1.0 if all include strings are present and exclude strings are absent, 0.0 otherwise
    """
    if result is None:
        return 0.0
    logger.info(f'Checking result: {result}, expected: {expected}')
    include = expected.get('include', [])
    exclude = expected.get('exclude', [])
    if all((r in result for r in include)) and all((r not in result for r in exclude)):
        return 1.0
    else:
        return 0.0

def check_min_extension_count__0f84311e(result_state, expected_state, **options):
    """
    Check if the number of unpacked extensions meets or exceeds the minimum count.

    Args:
        result_state: The actual count of unpacked extensions (int)
        expected_state: Dict containing 'min_count' key with the minimum required count
        **options: Additional options

    Returns:
        float: 1.0 if count >= min_count, 0.0 otherwise
    """
    if result_state is None:
        logger.error('No extension count data retrieved')
        return 0.0
    if not isinstance(expected_state, dict) or 'min_count' not in expected_state:
        logger.error("Expected state must be a dict with 'min_count' key")
        return 0.0
    min_count = expected_state['min_count']
    logger.info(f'Expected minimum count: {min_count}')
    logger.info(f'Actual unpacked extension count: {result_state}')
    if result_state >= min_count:
        logger.info(f'Extension count check passed: {result_state} >= {min_count}')
        return 1.0
    else:
        logger.info(f'Extension count check failed: {result_state} < {min_count}')
        return 0.0

def check_extension_contains__fc6800da1dfd116cace0d10a635a3df0(actual: str, expected: dict, **options) -> float:
    """
    Check if the VSCode extension name is contained in the installed extensions list.

    Args:
        actual (str): Output from 'code --list-extensions' command
        expected (dict): Expected rules with 'extension_substring' field
        **options: Additional options

    Returns:
        float: 1.0 if extension substring is found, 0.0 otherwise
    """
    if not actual:
        return 0.0
    extension_substring = expected.get('extension_substring', '')
    if not extension_substring:
        return 0.0
    if extension_substring.lower() in actual.lower():
        return 1.0
    return 0.0

def check_extension_description__ae6416e4(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Check if a specific extension has the expected description.

    Args:
        result: Dict with 'extensions' (mapping names to descriptions) and 'all_names' (list of extension names)
        expected: Rules dict with 'extension_name' and 'description' keys
        **options: Additional options

    Returns:
        1.0 if extension has the expected description, 0.0 otherwise
    """
    extension_name = expected.get('extension_name', '')
    expected_description = expected.get('description', '')
    if not extension_name or not expected_description:
        logger.warning('Missing extension_name or description in expected values')
        return 0.0
    extensions = result.get('extensions', {})
    all_names = result.get('all_names', [])
    logger.info(f"Checking for extension '{extension_name}' with description '{expected_description}'")
    logger.info(f'Found {len(all_names)} extensions: {all_names}')
    if extension_name in extensions and extensions[extension_name] == expected_description:
        logger.info(f"Extension '{extension_name}' found with correct description")
        return 1.0
    else:
        if extension_name not in extensions:
            logger.warning(f"Extension '{extension_name}' not found. Available extensions: {all_names}")
        else:
            logger.warning(f"Extension '{extension_name}' found but description mismatch. Expected: '{expected_description}', Got: '{extensions[extension_name]}'")
        return 0.0

def check_table_completeness__5a197d93(result: Dict[str, int], expected: Dict[str, Any], **options) -> float:
    """Check if table is sufficiently complete.

    Args:
        result: Dict with completeness statistics
        expected: Dict with minimum requirements for each field
        **options: Additional options

    Returns:
        Score between 0.0 and 1.0
    """
    min_rows = expected.get('min_rows', 1)
    min_addresses = expected.get('min_addresses', 1)
    min_contacts = expected.get('min_contacts', 1)
    min_websites = expected.get('min_websites', 1)
    score = 0.0
    if result['total_rows'] >= min_rows:
        score += 0.25
    else:
        score += 0.25 * (result['total_rows'] / min_rows)
    if result['filled_addresses'] >= min_addresses:
        score += 0.25
    else:
        score += 0.25 * (result['filled_addresses'] / min_addresses)
    if result['filled_contacts'] >= min_contacts:
        score += 0.25
    else:
        score += 0.25 * (result['filled_contacts'] / min_contacts)
    if result['filled_websites'] >= min_websites:
        score += 0.25
    else:
        score += 0.25 * (result['filled_websites'] / min_websites)
    return min(1.0, score)

def check_webext_manifest__d1a4f844666cd6560872c801fdefe60a(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if web extension manifest.json matches expected structure.

    Args:
        result: Manifest JSON dict from getter (or None if file not found)
        expected: Expected manifest structure with rules to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.warning('Manifest file not found')
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background', False)
    expected_has_content_scripts = expected.get('has_content_scripts', False)
    expected_has_page_action = expected.get('has_page_action', False)
    expected_description_empty = expected.get('description_empty', False)
    expected_no_browser_action = expected.get('no_browser_action', False)
    score = 0.0
    total_checks = 0
    if expected_name is not None:
        total_checks += 1
        if result.get('name') == expected_name:
            score += 1
            logger.info(f'Name matches: {expected_name}')
        else:
            logger.warning(f"Name mismatch: expected '{expected_name}', got '{result.get('name')}'")
    if expected_version is not None:
        total_checks += 1
        if result.get('version') == expected_version:
            score += 1
            logger.info(f'Version matches: {expected_version}')
        else:
            logger.warning(f"Version mismatch: expected '{expected_version}', got '{result.get('version')}'")
    if expected_has_background:
        total_checks += 1
        if 'background' in result and 'scripts' in result['background'] and (len(result['background']['scripts']) > 0):
            score += 1
            logger.info(f"Background scripts found: {result['background']['scripts']}")
        else:
            logger.warning('Background scripts not found')
    if expected_has_content_scripts:
        total_checks += 1
        if 'content_scripts' in result and len(result['content_scripts']) > 0:
            score += 1
            logger.info(f"Content scripts found: {len(result['content_scripts'])} entries")
        else:
            logger.warning('Content scripts not found')
    if expected_has_page_action:
        total_checks += 1
        if 'page_action' in result:
            score += 1
            logger.info('Page action found')
        else:
            logger.warning('Page action not found')
    if expected_description_empty:
        total_checks += 1
        description = result.get('description', '')
        if description in [None, '']:
            score += 1
            logger.info('Description is empty as expected')
        else:
            logger.warning(f"Description should be empty but got: '{description}'")
    if expected_no_browser_action:
        total_checks += 1
        if 'browser_action' not in result:
            score += 1
            logger.info('browser_action is absent as expected')
        else:
            logger.warning(f"browser_action should not be present but found: {result['browser_action']}")
    if total_checks == 0:
        logger.warning('No checks to perform')
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks} checks passed)')
    return final_score

def check_at_least_one_extension__33e136ec(result, expected, **options):
    """
    Check if at least one of the expected extensions is installed.

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' list of extension names
        **options: Additional options

    Returns:
        float: 1.0 if at least one expected extension is installed, 0.0 otherwise
    """
    if not result:
        result = []
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        return 0.0
    set_expected = set(expected_extensions)
    set_installed = set(result)
    matched = set_expected.intersection(set_installed)
    logger.info(f'Expected extensions (at least one): {expected_extensions}')
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Matched: {matched}')
    if len(matched) > 0:
        return 1.0
    else:
        return 0.0

def is_expected_bookmarks__a82b78bb_7fde_4cb3_94a4_035baf10bcf0_aug_14_task_verify_0(bookmarks: List[str], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.

    This is a fixed version that uses parameterized folder names from rule['names']
    instead of hardcoded 'Liked Authors'.
    """
    if not bookmarks:
        return 0.0
    elif rule['type'] == 'bookmark_bar_folders_names':
        bookmark_bar_folders_names = [bookmark['name'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder']
        return 1.0 if set(rule['names']).issubset(set(bookmark_bar_folders_names)) else 0.0
    elif rule['type'] == 'bookmark_bar_websites_urls':
        bookmark_bar_websites_urls = [bookmark['url'] for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'url']
        return 1.0 if set(rule['urls']).issubset(set(bookmark_bar_websites_urls)) else 0.0
    elif rule['type'] == 'bookmark_bar_websites_with_titles':
        bookmark_bar_websites = [{'url': bookmark['url'], 'name': bookmark['name']} for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'url']
        expected_bookmarks = rule['bookmarks']
        for expected in expected_bookmarks:
            expected_url = expected['url']
            expected_title = expected['title']
            found = False
            for actual in bookmark_bar_websites:
                if actual['url'] == expected_url and actual['name'] == expected_title:
                    found = True
                    break
            if not found:
                logger.info(f'Expected bookmark not found: URL={expected_url}, Title={expected_title}')
                return 0.0
        return 1.0
    elif rule['type'] == 'liked_authors_websites_urls':
        expected_folder_name = rule['names'][0] if rule.get('names') else 'Liked Authors'
        authors_folder = next((bookmark for bookmark in bookmarks['bookmark_bar']['children'] if bookmark['type'] == 'folder' and bookmark['name'] == expected_folder_name), None)
        if authors_folder:
            logger.info(f"'{expected_folder_name}' folder exists")
            authors_urls = [bookmark['url'] for bookmark in authors_folder['children'] if bookmark['type'] == 'url']
            logger.info(f"Here is the '{expected_folder_name}' folder's urls: {authors_urls}")
            urls = rule['urls']
            for (idx, url) in enumerate(urls):
                if isinstance(url, str):
                    urls[idx] = [url]
            combinations = product(*urls)
            for combination in combinations:
                if set(combination) == set(authors_urls):
                    return 1.0
            return 0.0
        else:
            logger.info(f"'{expected_folder_name}' folder not found")
            return 0.0
    else:
        raise TypeError(f"{rule['type']} not support yet!")

def check_chrome_startup_urls__715fc21b1c707dd79fc5ab9b6e4df514(result, expected, **options):
    """
    Check if Chrome startup URLs match the expected list.

    Args:
        result: List of current startup URLs from getter
        expected: Dictionary with 'urls' key containing expected URL list

    Returns:
        1.0 if URLs match expected list, 0.0 otherwise
    """
    expected_urls = expected.get('urls', [])
    logger.info(f'Current startup URLs: {result}')
    logger.info(f'Expected startup URLs: {expected_urls}')
    result_set = set(result) if result else set()
    expected_set = set(expected_urls) if expected_urls else set()
    if result_set == expected_set:
        return 1.0
    else:
        return 0.0

def check_webext_manifest__70dff70b667529be05282f3babd2fe6e(result: Optional[Dict[str, Any]], expected: Dict[str, Any], **options) -> float:
    """
    Check if web extension manifest.json matches expected structure.

    Args:
        result: Manifest JSON dict from getter (or None if file not found)
        expected: Expected manifest structure with rules to check
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if result is None:
        logger.warning('Manifest file not found')
        return 0.0
    expected_name = expected.get('name')
    expected_version = expected.get('version')
    expected_has_background = expected.get('has_background')
    expected_has_browser_action = expected.get('has_browser_action')
    expected_has_content_scripts = expected.get('has_content_scripts')
    score = 0.0
    total_checks = 0
    if expected_name is not None:
        total_checks += 1
        if result.get('name') == expected_name:
            score += 1
            logger.info(f'Name matches: {expected_name}')
        else:
            logger.warning(f"Name mismatch: expected '{expected_name}', got '{result.get('name')}'")
    if expected_version is not None:
        total_checks += 1
        if result.get('version') == expected_version:
            score += 1
            logger.info(f'Version matches: {expected_version}')
        else:
            logger.warning(f"Version mismatch: expected '{expected_version}', got '{result.get('version')}'")
    total_checks += 1
    description = result.get('description', '')
    if not description or description.strip() == '':
        score += 1
        logger.info('Description is blank as expected')
    else:
        logger.warning(f"Description should be blank but got: '{description}'")
    if expected_has_background is not None:
        total_checks += 1
        if expected_has_background:
            if 'background' in result and 'scripts' in result['background'] and (len(result['background']['scripts']) > 0):
                score += 1
                logger.info(f"Background scripts found: {result['background']['scripts']}")
            else:
                logger.warning('Background scripts not found')
        elif 'background' not in result or 'scripts' not in result.get('background', {}) or len(result.get('background', {}).get('scripts', [])) == 0:
            score += 1
            logger.info('Background scripts correctly absent')
        else:
            logger.warning('Background scripts should be absent but were found')
    if expected_has_browser_action is not None:
        total_checks += 1
        if expected_has_browser_action:
            if 'browser_action' in result:
                score += 1
                logger.info('Browser action found')
            else:
                logger.warning('Browser action not found')
        elif 'browser_action' not in result:
            score += 1
            logger.info('Browser action correctly absent')
        else:
            logger.warning('Browser action should be absent but was found')
    if expected_has_content_scripts is not None:
        total_checks += 1
        if expected_has_content_scripts:
            if 'content_scripts' in result and len(result['content_scripts']) > 0:
                score += 1
                logger.info(f"Content scripts found: {len(result['content_scripts'])} entries")
            else:
                logger.warning('Content scripts not found')
        elif 'content_scripts' not in result or len(result.get('content_scripts', [])) == 0:
            score += 1
            logger.info('Content scripts correctly absent')
        else:
            logger.warning(f"Content scripts should be absent but found {len(result['content_scripts'])} entries")
    if total_checks == 0:
        logger.warning('No checks to perform')
        return 0.0
    final_score = score / total_checks
    logger.info(f'Final score: {final_score} ({score}/{total_checks} checks passed)')
    return final_score

def check_gdrive_pdf_in_folder__42b1e128(result, expected, **options):
    """
    Check if a PDF file exists in the correct Google Drive folder.

    Args:
        result: Dict from getter with file_exists, folder_exists, in_folder
        expected: Dict with file_should_exist and in_folder expectations
        **options: Additional options

    Returns:
        float: 1.0 if file exists in folder, 0.0 otherwise
    """
    if not isinstance(result, dict):
        logger.error(f'Result is not a dict: {result}')
        return 0.0
    if not isinstance(expected, dict):
        logger.error(f'Expected is not a dict: {expected}')
        return 0.0
    file_should_exist = expected.get('file_should_exist', True)
    should_be_in_folder = expected.get('in_folder', True)
    file_exists = result.get('file_exists', False)
    in_folder = result.get('in_folder', False)
    logger.info(f'File exists: {file_exists}, In folder: {in_folder}')
    logger.info(f'Expected - File should exist: {file_should_exist}, Should be in folder: {should_be_in_folder}')
    if file_should_exist and (not file_exists):
        logger.info('❌ File does not exist')
        return 0.0
    if should_be_in_folder and (not in_folder):
        logger.info('❌ File exists but not in the correct folder')
        return 0.0
    logger.info('✅ File exists in the correct folder')
    return 1.0

def is_extension_installed__2dcff8b4(actual: str, expected: Dict, **options):
    """Check if a VS Code extension is installed.

    Args:
        actual: Output from 'code --list-extensions' command
        expected: Dict with 'type' ('contain' or 'not_contain') and 'expected' (extension ID)

    Returns:
        float: 1.0 if check passes, 0.0 otherwise
    """
    if expected['type'] == 'contain':
        if expected['expected'] in actual:
            return 1.0
        return 0.0
    elif expected['type'] == 'not_contain':
        if expected['expected'] not in actual:
            return 1.0
        return 0.0
    else:
        raise NotImplementedError(f"Unknown type: {expected['type']}")

def check_gdrive_files_exist__74b11cf6(result_state: Union[List[str], str, None], expected_state: Dict[str, Any], **options) -> float:
    """
    Check if Google Drive files were successfully downloaded, are in .eml format,
    and contain the expected email content from Thunderbird Bills folder.

    This evaluator verifies:
    1. File count matches expected
    2. All files have .eml extension
    3. Subject lines were preserved as filenames (with filesystem-safe character substitution)
    4. Email content contains expected markers (sender addresses, subject patterns)

    Args:
        result_state: List of local file paths returned by get_googledrive_file, or None if download failed
        expected_state: Dict containing verification criteria:
            - file_count: Expected number of files
            - expected_emails: List of dicts with email metadata (subject, sender patterns)
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0 (1.0 if all files exist and match expected emails, 0.0 otherwise)
    """
    if result_state is None:
        return 0.0
    expected_count = expected_state.get('file_count', 0)
    expected_emails = expected_state.get('expected_emails', [])
    if isinstance(result_state, str):
        result_state = [result_state]
    if not isinstance(result_state, list):
        return 0.0
    downloaded_files = [path for path in result_state if path is not None]
    existing_files = [path for path in downloaded_files if os.path.exists(path)]
    if len(existing_files) != expected_count:
        return 0.0
    for file_path in existing_files:
        if not file_path.endswith('.eml'):
            return 0.0
    if expected_emails:
        verified_emails = []
        for file_path in existing_files:
            try:
                with open(file_path, 'rb') as f:
                    msg = email.message_from_binary_file(f, policy=policy.default)
                subject = msg.get('Subject', '')
                sender = msg.get('From', '')
                matched = False
                for expected_email in expected_emails:
                    expected_subject = expected_email.get('subject', '')
                    expected_sender_pattern = expected_email.get('sender_pattern', '')
                    if subject.strip().lower() == expected_subject.strip().lower():
                        if expected_sender_pattern:
                            if re.search(expected_sender_pattern, sender, re.IGNORECASE):
                                matched = True
                                verified_emails.append(expected_email)
                                break
                        else:
                            matched = True
                            verified_emails.append(expected_email)
                            break
                if not matched:
                    return 0.0
            except Exception as e:
                return 0.0
        if len(verified_emails) != len(expected_emails):
            return 0.0
    return 1.0

def check_extension_source_type__ae6416e4(result: Dict[str, str], expected: Dict[str, Any], **options) -> float:
    """Check if a specific extension has the expected source type.

    Args:
        result: Dict mapping extension names to source types
        expected: Rules dict with 'extension_name' and 'source_type' keys
        **options: Additional options

    Returns:
        1.0 if extension has the expected source type, 0.0 otherwise
    """
    extension_name = expected.get('extension_name', '')
    expected_source = expected.get('source_type', '')
    if not extension_name or not expected_source:
        return 0.0
    if extension_name in result and result[extension_name] == expected_source:
        return 1.0
    else:
        return 0.0

def check_chrome_restore_setting__9cb5a6acff7768f53f96a10b1f86ac95(result, expected, **options):
    """
    Check if Chrome restore_on_startup setting matches expected value.

    Args:
        result: Current restore_on_startup value from getter
        expected: Dictionary with 'setting_value' key

    Returns:
        1.0 if setting matches expected, 0.0 otherwise
    """
    expected_value = expected.get('setting_value', 5)
    logger.info(f'Current restore setting: {result}')
    logger.info(f'Expected restore setting: {expected_value}')
    if result == expected_value:
        return 1.0
    else:
        return 0.0

def check_chrome_ext_subset__d7f8f8dc6935bc142e9b5dc629fdadfe(result: List[str], expected: Dict, **options) -> float:
    """Check if expected Chrome extensions are installed (subset check).

    Args:
        result: List of installed extension names from getter
        expected: Dict with 'expected' key containing list of required extension names
        **options: Additional options (not used)

    Returns:
        1.0 if all expected extensions are installed, 0.0 otherwise
    """
    if not result:
        logger.info('No extensions found in result')
        return 0.0
    expected_extensions = expected.get('expected', [])
    if not expected_extensions:
        logger.warning('No expected extensions specified')
        return 0.0
    logger.info(f'Installed extensions: {result}')
    logger.info(f'Expected extensions: {expected_extensions}')
    set_expected = set(expected_extensions)
    set_installed = set(result)
    if set_expected.issubset(set_installed):
        logger.info('All expected extensions are installed')
        return 1.0
    else:
        missing = set_expected - set_installed
        logger.info(f'Missing extensions: {missing}')
        return 0.0

def check_git_url_contains__07fbe1c4(result, expected, **options):
    """Check if git URL contains expected substring.

    Args:
        result: Git remote URL string
        expected: Expected substring in URL
        **options: Additional options

    Returns:
        float: 1.0 if expected substring is in URL, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    if expected in result:
        return 1.0
    return 0.0

def check_table_left_aligned__b943f0fd(result, expected, **options):
    """Check if table is aligned to the left margin.

    Args:
        result: Table position dict from getter
        expected: Dict with 'max_left' threshold
        **options: Additional options

    Returns:
        float: 1.0 if table left < max_left, 0.0 otherwise
    """
    if result is None:
        return 0.0
    max_left = expected.get('max_left', 0)
    if result['left'] < max_left:
        return 1.0
    else:
        return 0.0

def is_expected_bookmarks__2ad9387a(bookmarks: Dict[str, Any], rule: Dict[str, Any]) -> float:
    """
    Checks if the expected bookmarks are in Chrome.

    Supports bookmark_bar_websites_with_titles type to verify both URL and title.

    Args:
        bookmarks: Chrome bookmarks data structure
        rule: Rule dict with 'type' and 'bookmarks' keys
              For type='bookmark_bar_websites_with_titles':
                - bookmarks: List of dicts with 'url' and 'title' keys

    Returns:
        float: 1.0 if all expected bookmarks found with correct URL and title, 0.0 otherwise
    """
    if not bookmarks:
        logger.info('No bookmarks data provided')
        return 0.0
    if rule['type'] == 'bookmark_bar_websites_with_titles':
        if 'bookmark_bar' not in bookmarks:
            logger.info('No bookmark_bar found in bookmarks')
            return 0.0
        if 'children' not in bookmarks['bookmark_bar']:
            logger.info('No children in bookmark_bar')
            return 0.0
        bookmark_bar_websites = []
        for bookmark in bookmarks['bookmark_bar']['children']:
            if bookmark.get('type') == 'url':
                bookmark_bar_websites.append({'url': bookmark.get('url', ''), 'title': bookmark.get('name', '')})
        logger.info(f'Found {len(bookmark_bar_websites)} bookmarks in bookmark bar: {bookmark_bar_websites}')
        expected_bookmarks = rule.get('bookmarks', [])
        for expected in expected_bookmarks:
            expected_url = expected.get('url', '')
            expected_title = expected.get('title', '')
            expected_url_normalized = expected_url.rstrip('/')
            found = False
            for actual in bookmark_bar_websites:
                actual_url_normalized = actual['url'].rstrip('/')
                if actual_url_normalized == expected_url_normalized and actual['title'] == expected_title:
                    found = True
                    logger.info(f'Found matching bookmark: URL={expected_url}, Title={expected_title}')
                    break
            if not found:
                logger.info(f'Expected bookmark not found: URL={expected_url}, Title={expected_title}')
                return 0.0
        return 1.0
    else:
        raise TypeError(f"{rule['type']} not supported by is_expected_bookmarks__2ad9387a")

def check_gdrive_files__59579682(result, expected, **options):
    """Check if expected files are present in Google Drive.

    Args:
        result: List of file names from getter
        expected: Expected file list configuration
        **options: Additional options

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, list):
        return 0.0
    expected_files = expected.get('files', [])
    if not expected_files:
        return 0.0
    found_count = 0
    for expected_file in expected_files:
        if expected_file in result:
            found_count += 1
    return found_count / len(expected_files)

def check_chrome_extension_manifest__1e834a0e2fee4e317d31e5d8fca95c5c(result: Dict[str, Any], expected: Dict[str, Any], **options) -> float:
    """Validate Chrome extension manifest and files.

    Args:
        result: Dict from getter with 'manifest', 'files_exist', 'all_files_exist'
        expected: Dict with expected manifest structure rules:
            - manifest_checks: List of dicts with 'key' (list of nested keys) and 'value' (expected value)
        **options: Additional options (not used)

    Returns:
        float: Score between 0.0 and 1.0
    """
    if not isinstance(result, dict) or not isinstance(expected, dict):
        return 0.0
    manifest = result.get('manifest', {})
    all_files_exist = result.get('all_files_exist', False)
    file_score = 0.5 if all_files_exist else 0.0
    manifest_checks = expected.get('manifest_checks', [])
    if not manifest_checks:
        return file_score
    passed_checks = 0
    total_checks = len(manifest_checks)
    for check in manifest_checks:
        key_path = check.get('key', [])
        expected_value = check.get('value')
        current = manifest
        try:
            for key in key_path:
                current = current[key]
            if current == expected_value:
                passed_checks += 1
        except (KeyError, TypeError):
            pass
    manifest_score = 0.5 * (passed_checks / total_checks) if total_checks > 0 else 0.0
    return file_score + manifest_score

def check_table_right_aligned__d963ccf8(result, expected, **options):
    """Check if table is aligned to the right margin.

    Args:
        result: Table position dict from getter
        expected: Dict with 'min_left' threshold
        **options: Additional options

    Returns:
        float: 1.0 if table left > min_left, 0.0 otherwise
    """
    if result is None:
        return 0.0
    min_left = expected.get('min_left', 0)
    if result['left'] > min_left:
        return 1.0
    else:
        return 0.0

def check_block_third_party_cookies__d0393d13df595b6db99860dc4f30ea7b(cookie_settings, rule):
    """
    Check if the third-party cookie blocking setting is as expected.
    Args:
        cookie_settings: Dict containing block_third_party_cookies from Chrome preferences
        rule: Dict with validation rules (enabled: bool)
    Returns:
        float: 1.0 if matches expected state, 0.0 otherwise
    """
    is_blocked = cookie_settings.get('block_third_party_cookies', False)
    expected_blocked = rule.get('enabled', True)
    return 1.0 if is_blocked == expected_blocked else 0.0

def check_chrome_enhanced_safe_browsing__a6a55567(result, expected, **options):
    """Check if Chrome Enhanced Safe Browsing is in expected state.

    Args:
        result: Enhanced safe browsing settings dict from getter
        expected: Expected rules dict with 'enabled' key
        **options: Additional options

    Returns:
        float: 1.0 if matches, 0.0 otherwise
    """
    if not result or not expected:
        return 0.0
    expected_enabled = expected.get('enabled', False)
    actual_enabled = result.get('enabled', False)
    logger.info(f'Expected enhanced safe browsing enabled: {expected_enabled}')
    logger.info(f'Actual enhanced safe browsing enabled: {actual_enabled}')
    return 1.0 if actual_enabled == expected_enabled else 0.0
