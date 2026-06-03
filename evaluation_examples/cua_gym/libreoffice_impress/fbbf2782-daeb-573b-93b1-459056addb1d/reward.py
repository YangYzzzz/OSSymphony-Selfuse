"""
Reward Script: Build custom slide layouts and apply them to specific slides
Task ID: impress_gf4_012
Domain: libreoffice_impress
Scoring:
  Component 1: Four custom layouts defined in master (0.30)
  Component 2: Slide 1 uses custom 'Title Only' layout (not default) (0.15)
  Component 3: Slides 5 and 10 use 'Two Column' layout (0.20)
  Component 4: Slide 15 uses 'Section Break' layout (0.15)
  Component 5: Section Break layout has dark solid background (0.10)
  Component 6: Non-target slides retain default layout (0.10)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_012'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_layout_index(slide, master):
    """Find the index of a slide's layout within its master."""
    layout = slide.slide_layout
    for i, l in enumerate(master.slide_layouts):
        if l == layout:
            return i
    return -1


def get_custom_layout_names(master, default_count=11):
    """Get names of custom layouts (those beyond the default set)."""
    names = []
    for i, layout in enumerate(master.slide_layouts):
        if i >= default_count:
            names.append(layout.name)
    return names


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: 16 slides
    if len(prs.slides) != 16:
        print(f"PRECONDITION FAIL: Expected 16 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]
    total_layouts = len(master.slide_layouts)

    # Component 1: Four custom layouts defined (0.30 points)
    # The initial file has 11 default layouts. Golden should have 15 (11 default + 4 custom).
    # We check for the existence of custom layouts named:
    # 'Title Only' (custom, beyond index 10), 'Two Column', 'Quote Layout', 'Section Break'
    try:
        required_custom_names = {'Two Column', 'Quote Layout', 'Section Break'}
        custom_names = get_custom_layout_names(master, default_count=11)
        custom_names_set = set(custom_names)

        # Check that we have at least 4 custom layouts (beyond the 11 defaults)
        has_enough_custom = len(custom_names) >= 4

        # Check that the required named layouts exist among custom layouts
        # 'Title Only' is special - it shares a name with the default layout 5, but must be a
        # SEPARATE custom layout at index >= 11
        has_custom_title_only = 'Title Only' in custom_names
        has_required_names = required_custom_names.issubset(custom_names_set)

        if has_enough_custom and has_custom_title_only and has_required_names:
            print(f"PASS: Component 1 -- 4 custom layouts found: {custom_names} (0.30 pts)")
            total_score += 0.30
        else:
            missing = required_custom_names - custom_names_set
            print(f"FAIL: Component 1 -- Custom layouts: {custom_names}, "
                  f"enough={has_enough_custom}, custom_title_only={has_custom_title_only}, "
                  f"missing_required={missing}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 1 uses custom 'Title Only' layout (index >= 11) (0.15 points)
    # In the initial file, slide 1 uses default 'Title Only' at index 5.
    # In the golden file, it should use the CUSTOM 'Title Only' at index >= 11.
    try:
        slide1 = prs.slides[0]
        layout_idx = get_layout_index(slide1, master)
        layout_name = slide1.slide_layout.name

        if layout_name == 'Title Only' and layout_idx >= 11:
            print(f"PASS: Component 2 -- Slide 1 uses custom 'Title Only' "
                  f"(layout index {layout_idx}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Slide 1 layout: name='{layout_name}', "
                  f"index={layout_idx} (expected custom 'Title Only' at index >= 11)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slides 5 and 10 use 'Two Column' layout (0.20 points)
    # Both must use it. Award 0.10 per slide.
    try:
        slide5 = prs.slides[4]
        slide10 = prs.slides[9]
        slide5_name = slide5.slide_layout.name
        slide10_name = slide10.slide_layout.name

        c3_score = 0.0
        if slide5_name == 'Two Column':
            print(f"PASS: Component 3a -- Slide 5 uses 'Two Column' layout (0.10 pts)")
            c3_score += 0.10
        else:
            print(f"FAIL: Component 3a -- Slide 5 layout: '{slide5_name}' (expected 'Two Column')")

        if slide10_name == 'Two Column':
            print(f"PASS: Component 3b -- Slide 10 uses 'Two Column' layout (0.10 pts)")
            c3_score += 0.10
        else:
            print(f"FAIL: Component 3b -- Slide 10 layout: '{slide10_name}' (expected 'Two Column')")

        if c3_score > 0:
            total_score += c3_score
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 15 uses 'Section Break' layout (0.15 points)
    try:
        slide15 = prs.slides[14]
        slide15_name = slide15.slide_layout.name

        if slide15_name == 'Section Break':
            print(f"PASS: Component 4 -- Slide 15 uses 'Section Break' layout (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Slide 15 layout: '{slide15_name}' "
                  f"(expected 'Section Break')")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Section Break layout has dark solid background (0.10 points)
    # The layout or the slide itself should have a solid dark background.
    # Dark = R,G,B each <= 80 (out of 255).
    try:
        slide15 = prs.slides[14]
        bg_fill = slide15.background.fill
        dark_bg_found = False
        bg_color_str = "none"

        if bg_fill.type == 1:  # SOLID
            rgb = bg_fill.fore_color.rgb
            bg_color_str = str(rgb)
            r, g, b = int(str(rgb)[:2], 16), int(str(rgb)[2:4], 16), int(str(rgb)[4:6], 16)
            if r <= 80 and g <= 80 and b <= 80:
                dark_bg_found = True

        # Also check via layout if slide inherits
        if not dark_bg_found:
            layout_fill = slide15.slide_layout.background.fill
            if layout_fill.type == 1:
                rgb = layout_fill.fore_color.rgb
                bg_color_str = str(rgb)
                r, g, b = int(str(rgb)[:2], 16), int(str(rgb)[2:4], 16), int(str(rgb)[4:6], 16)
                if r <= 80 and g <= 80 and b <= 80:
                    dark_bg_found = True

        if dark_bg_found:
            print(f"PASS: Component 5 -- Section Break has dark background "
                  f"(color: {bg_color_str}) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 -- Slide 15 background not dark "
                  f"(fill type: {bg_fill.type}, color: {bg_color_str})")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Non-target slides retain default layout (0.10 points)
    # Slides not in {1, 5, 10, 15} should still use the default 'Title Only' layout at index 5.
    # This is a preservation check - ensures golden_patch didn't accidentally change other slides.
    try:
        target_slides = {0, 4, 9, 14}  # 0-indexed
        non_target_correct = 0
        non_target_total = 16 - len(target_slides)  # 12 slides

        for i, slide in enumerate(prs.slides):
            if i in target_slides:
                continue
            layout_idx = get_layout_index(slide, master)
            # Default 'Title Only' is at index 5
            if layout_idx == 5:
                non_target_correct += 1
            else:
                print(f"  NOTE: Slide {i+1} has non-default layout index {layout_idx} "
                      f"(name: {slide.slide_layout.name})")

        if non_target_correct == non_target_total:
            print(f"PASS: Component 6 -- All {non_target_total} non-target slides "
                  f"retain default layout (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 -- {non_target_correct}/{non_target_total} "
                  f"non-target slides have default layout")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state("libreoffice_impress")

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
