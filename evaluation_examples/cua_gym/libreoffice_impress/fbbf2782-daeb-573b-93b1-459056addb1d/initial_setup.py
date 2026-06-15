"""
Initial Setup: Build a corporate pitch presentation with 16 slides
Task ID: impress_gf4_012
Domain: libreoffice_impress

Creates a 16-slide Corporate_Pitch.pptx with realistic content.
All slides use the default layout. No custom layouts are defined.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet_points(text_frame, items, font_size=Pt(16), color=None):
    """Add multiple bullet points to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = font_size
            if color:
                run.font.color.rgb = color


def create_initial():
    prs = Presentation()

    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # We'll use layout index 5 (Blank) for all slides and add content manually
    # This gives us a clean single-layout presentation
    blank_layout = prs.slide_layouts[5]

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    add_text_to_shape(txBox, "NovaTech Solutions", font_size=Pt(44), bold=True,
                      color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.8), Inches(9), Inches(1))
    add_text_to_shape(txBox2, "Series B Fundraising Pitch — Q2 2025",
                      font_size=Pt(24), color=RGBColor(0x5A, 0x5A, 0x5A), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Agenda", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "1. Company Overview & Mission",
        "2. Market Opportunity & TAM Analysis",
        "3. Product Demo & Technology Stack",
        "4. Revenue Model & Unit Economics",
        "5. Go-to-Market Strategy",
        "6. Competitive Landscape",
        "7. Traction & Key Metrics",
        "8. Financial Projections (2025-2028)",
        "9. Team & Advisory Board",
        "10. Investment Ask & Use of Funds",
    ]
    add_bullet_points(tf, items, font_size=Pt(20))

    # --- Slide 3: Company Overview ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Company Overview", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Founded in 2021 by former AWS and Stripe engineers",
        "Headquartered in San Francisco with offices in London and Singapore",
        "85 employees across engineering, sales, and operations",
        "Mission: Democratize enterprise-grade AI infrastructure for mid-market companies",
        "Core product: NovaPlatform — a unified AI deployment and monitoring suite",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 4: The Problem ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "The Problem We Solve", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Mid-market companies spend $2.3M annually on fragmented AI tooling",
        "Average deployment time for a single ML model: 14 weeks",
        "78% of AI projects in mid-market never reach production",
        "Existing solutions are either too expensive (Databricks) or too complex (custom Kubernetes)",
        "No unified platform for model deployment, monitoring, and governance",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 5: Market Opportunity ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Market Opportunity", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Total Addressable Market: $47B by 2027",
        "Serviceable Addressable Market: $12B",
        "Serviceable Obtainable Market: $1.8B",
        "Growing at 34% CAGR (Gartner, 2024)",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))
    txBox3 = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5), Inches(4.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    items3 = [
        "Key Tailwinds:",
        "  - Regulatory push for AI governance",
        "  - Explosion of open-source LLMs",
        "  - Mid-market digital transformation",
        "  - Talent shortage in MLOps",
    ]
    add_bullet_points(tf3, items3, font_size=Pt(18))

    # --- Slide 6: Product Overview ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "NovaPlatform: Product Overview", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "One-Click Model Deployment — deploy any model to production in under 2 hours",
        "Real-Time Monitoring Dashboard — track latency, drift, and accuracy metrics",
        "Automated A/B Testing — built-in experimentation framework",
        "Compliance & Governance Module — SOC 2 and GDPR compliant audit trails",
        "Integrated Cost Optimizer — reduce cloud spend by up to 40%",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 7: Technology Stack ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Technology Architecture", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Backend: Go microservices on Kubernetes with gRPC",
        "ML Runtime: Custom container orchestration with GPU autoscaling",
        "Frontend: React + TypeScript with real-time WebSocket updates",
        "Data Layer: PostgreSQL + ClickHouse for analytics, Redis for caching",
        "Infrastructure: Multi-cloud (AWS, GCP, Azure) with Terraform IaC",
        "Security: Zero-trust architecture, end-to-end encryption at rest and in transit",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 8: Revenue Model ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Revenue Model & Unit Economics", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "SaaS subscription model with three tiers:",
        "  Starter: $2,500/month (up to 5 models, 10 users)",
        "  Growth: $8,000/month (up to 25 models, 50 users)",
        "  Enterprise: $25,000/month (unlimited, dedicated support)",
        "Average Contract Value (ACV): $96,000",
        "Gross Margin: 78%",
        "Net Dollar Retention: 135%",
        "CAC Payback Period: 11 months",
    ]
    add_bullet_points(tf, items, font_size=Pt(16))

    # --- Slide 9: Go-to-Market ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Go-to-Market Strategy", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Product-led growth: Free tier for individual developers",
        "Enterprise sales team: 12 AEs targeting mid-market (500-5000 employees)",
        "Channel partnerships: AWS Marketplace, Snowflake Partner Connect",
        "Content marketing: Technical blog with 45K monthly readers",
        "Community: Open-source SDK with 8,200 GitHub stars",
        "Events: Sponsoring MLConf, AI Summit, and KubeCon 2025",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 10: Competitive Landscape ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Competitive Landscape", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Enterprise Players:",
        "  - Databricks ($43B valuation, complex)",
        "  - AWS SageMaker (vendor lock-in)",
        "  - Google Vertex AI (GCP-only)",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))
    txBox3 = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(5), Inches(4.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    items3 = [
        "Our Differentiation:",
        "  - Multi-cloud native from day one",
        "  - 10x faster deployment vs competitors",
        "  - Purpose-built for mid-market budget",
        "  - Integrated governance out of the box",
    ]
    add_bullet_points(tf3, items3, font_size=Pt(18))

    # --- Slide 11: Traction & Metrics ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Traction & Key Metrics", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "ARR: $4.2M (up from $1.1M 12 months ago — 282% YoY growth)",
        "Customers: 47 paying accounts (up from 14)",
        "Logo churn: 3% annually",
        "NPS Score: 72",
        "Models deployed on platform: 1,240+",
        "API calls processed monthly: 890M",
        "Notable customers: Rivian, Warby Parker, DoorDash, Gusto",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 12: Financial Projections ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Financial Projections (2025-2028)", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))

    # Add a simple table
    table_shape = slide.shapes.add_table(5, 5, Inches(1), Inches(2), Inches(11), Inches(3))
    table = table_shape.table
    headers = ["Metric", "2025", "2026", "2027", "2028"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        for run in table.cell(0, i).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data_rows = [
        ["ARR ($M)", "$8.5", "$22.0", "$48.0", "$95.0"],
        ["Customers", "110", "280", "520", "900"],
        ["Employees", "130", "220", "380", "550"],
        ["Gross Margin", "79%", "81%", "83%", "85%"],
    ]
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val
            for run in table.cell(r, c).text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # --- Slide 13: Team ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Leadership Team", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Elena Rodriguez, CEO — Ex-VP Engineering at Stripe, Stanford CS",
        "James Park, CTO — Former Principal Engineer at AWS, MIT PhD",
        "Priya Sharma, VP Product — Previously at Datadog, 12 years in MLOps",
        "Marcus Williams, VP Sales — Built enterprise sales at Snowflake ($0-$50M)",
        "Aisha Okonkwo, VP Engineering — Former Tech Lead at Google Brain",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 14: Advisory Board ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Advisory Board", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Dr. Andrew Chen — Professor of ML at Carnegie Mellon, former Uber Chief Scientist",
        "Sarah Liu — Managing Director at Sequoia Capital, board member at Scale AI",
        "Robert Nakamura — Former CIO at Goldman Sachs, enterprise governance expert",
        "Lisa Thompson — CEO of CloudSecure (acquired by Palo Alto Networks, 2023)",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 15: Investment Ask ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    add_text_to_shape(txBox, "Investment Ask", font_size=Pt(36), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C))
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(10), Inches(4.5))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Raising $35M Series B at $175M pre-money valuation",
        "Use of Funds:",
        "  45% — Engineering & Product (hire 50 engineers)",
        "  25% — Sales & Marketing (expand enterprise team)",
        "  15% — Infrastructure (multi-cloud expansion)",
        "  10% — International expansion (EMEA, APAC)",
        "  5%  — G&A and working capital",
    ]
    add_bullet_points(tf, items, font_size=Pt(18))

    # --- Slide 16: Thank You / Contact ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(1.5))
    add_text_to_shape(txBox, "Thank You", font_size=Pt(44), bold=True,
                      color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(2))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = [
        "Elena Rodriguez, CEO",
        "elena@novatech.io | +1 (415) 555-0192",
        "www.novatech.io",
    ]
    add_bullet_points(tf, items, font_size=Pt(20), color=RGBColor(0x5A, 0x5A, 0x5A))
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
