"""
Initial Setup: Edit master slide body text style
Task ID: impress_ma_014
Domain: libreoffice_impress

Creates a 30-slide lecture presentation with Liberation Sans 18pt body text
and single line spacing on the slide master.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Lecture topics for 30 slides
LECTURE_TOPICS = [
    ("Introduction to Machine Learning", "Course Overview and Prerequisites"),
    ("Supervised Learning Foundations", "Key Concepts in Classification and Regression"),
    ("Linear Regression", "Modeling Continuous Outcomes"),
    ("Logistic Regression", "Binary and Multiclass Classification"),
    ("Decision Trees", "Recursive Partitioning Methods"),
    ("Random Forests", "Ensemble Learning with Bagging"),
    ("Support Vector Machines", "Maximum Margin Classifiers"),
    ("Neural Network Basics", "Perceptrons and Activation Functions"),
    ("Deep Learning Architectures", "CNNs, RNNs, and Transformers"),
    ("Backpropagation", "Gradient Computation and Optimization"),
    ("Regularization Techniques", "L1, L2, and Dropout Methods"),
    ("Model Evaluation Metrics", "Precision, Recall, and F1-Score"),
    ("Cross-Validation", "K-Fold and Stratified Sampling"),
    ("Hyperparameter Tuning", "Grid Search and Bayesian Optimization"),
    ("Unsupervised Learning", "Clustering and Dimensionality Reduction"),
    ("K-Means Clustering", "Partitioning Data into Groups"),
    ("Principal Component Analysis", "Dimensionality Reduction Techniques"),
    ("Natural Language Processing", "Text Processing and Embeddings"),
    ("Computer Vision", "Image Recognition and Object Detection"),
    ("Reinforcement Learning", "Agents, Rewards, and Policies"),
    ("Generative Models", "GANs and Variational Autoencoders"),
    ("Transfer Learning", "Fine-Tuning Pretrained Models"),
    ("Ethics in AI", "Bias, Fairness, and Accountability"),
    ("Data Preprocessing", "Cleaning, Normalization, and Feature Engineering"),
    ("Feature Selection", "Filter, Wrapper, and Embedded Methods"),
    ("Time Series Analysis", "Forecasting and Temporal Patterns"),
    ("Recommendation Systems", "Collaborative and Content-Based Filtering"),
    ("Deploying ML Models", "Serving, Monitoring, and MLOps"),
    ("Research Frontiers", "Current Trends and Open Problems"),
    ("Course Summary", "Key Takeaways and Next Steps"),
]

BODY_CONTENTS = [
    "This course covers the fundamental concepts of machine learning, from basic algorithms to advanced deep learning techniques. Students will learn to build, evaluate, and deploy ML models.\n\nPrerequisites include linear algebra, probability, and Python programming.",
    "Supervised learning uses labeled data to train predictive models. The two main categories are classification (discrete labels) and regression (continuous values).\n\nKey concepts: training set, test set, overfitting, underfitting, bias-variance tradeoff.",
    "Linear regression models the relationship between features and a continuous target variable using a linear function.\n\nTopics: ordinary least squares, gradient descent, multiple regression, polynomial features, residual analysis.",
    "Logistic regression extends linear models to classification by applying the sigmoid function.\n\nTopics: log-odds, decision boundary, softmax for multiclass, regularized logistic regression.",
    "Decision trees partition the feature space into rectangular regions using a series of binary splits.\n\nTopics: information gain, Gini impurity, pruning strategies, handling missing values.",
    "Random forests combine multiple decision trees trained on bootstrap samples to reduce overfitting.\n\nTopics: bagging, feature randomization, out-of-bag error, variable importance measures.",
    "SVMs find the hyperplane that maximizes the margin between classes in the feature space.\n\nTopics: kernel trick, soft margin, RBF kernel, support vectors, dual formulation.",
    "Neural networks consist of layers of interconnected neurons that learn hierarchical feature representations.\n\nTopics: perceptron, sigmoid/ReLU activation, feedforward networks, universal approximation theorem.",
    "Modern deep learning architectures address specific data modalities and sequence processing needs.\n\nTopics: convolutional layers, recurrent connections, attention mechanisms, transformer architecture.",
    "Backpropagation computes gradients of the loss function with respect to each weight using the chain rule.\n\nTopics: computational graphs, vanishing gradients, gradient clipping, learning rate schedules.",
    "Regularization prevents overfitting by adding constraints or noise during training.\n\nTopics: L1 (Lasso), L2 (Ridge), elastic net, dropout, early stopping, data augmentation.",
    "Model evaluation requires appropriate metrics that align with the business objective and data characteristics.\n\nTopics: accuracy, precision, recall, F1-score, ROC-AUC, confusion matrix, calibration.",
    "Cross-validation provides robust estimates of model performance by training and testing on different data splits.\n\nTopics: k-fold CV, stratified CV, leave-one-out, nested CV for hyperparameter selection.",
    "Hyperparameter tuning optimizes model configuration parameters that are not learned during training.\n\nTopics: grid search, random search, Bayesian optimization, Hyperband, early stopping criteria.",
    "Unsupervised learning discovers patterns in unlabeled data without explicit target variables.\n\nTopics: clustering, dimensionality reduction, anomaly detection, density estimation.",
    "K-means partitions n observations into k clusters by minimizing within-cluster sum of squares.\n\nTopics: initialization methods, elbow method, silhouette score, k-means++ algorithm.",
    "PCA projects data onto orthogonal components that capture maximum variance in the original features.\n\nTopics: eigenvalue decomposition, scree plot, explained variance ratio, kernel PCA.",
    "NLP enables machines to understand, interpret, and generate human language from text data.\n\nTopics: tokenization, word embeddings, TF-IDF, BERT, GPT, sentiment analysis.",
    "Computer vision extracts meaningful information from images and video using deep learning models.\n\nTopics: convolution, pooling, object detection (YOLO, R-CNN), image segmentation, data augmentation.",
    "Reinforcement learning trains agents to make sequential decisions by maximizing cumulative rewards.\n\nTopics: Markov decision processes, Q-learning, policy gradient, exploration vs exploitation.",
    "Generative models learn the underlying data distribution to create new, realistic samples.\n\nTopics: GANs (generator/discriminator), VAEs (encoder/decoder), diffusion models, evaluation metrics.",
    "Transfer learning leverages knowledge from pretrained models to improve performance on new tasks.\n\nTopics: feature extraction, fine-tuning, domain adaptation, few-shot learning.",
    "Responsible AI development requires careful consideration of societal impact and ethical implications.\n\nTopics: algorithmic bias, fairness metrics, explainability, privacy, accountability frameworks.",
    "Data preprocessing transforms raw data into a clean, structured format suitable for modeling.\n\nTopics: missing value imputation, outlier detection, normalization, encoding categorical variables.",
    "Feature selection identifies the most relevant variables to improve model performance and interpretability.\n\nTopics: correlation analysis, mutual information, recursive feature elimination, LASSO.",
    "Time series analysis models temporal dependencies and trends in sequential data for forecasting.\n\nTopics: ARIMA, seasonal decomposition, exponential smoothing, LSTM for sequences.",
    "Recommendation systems predict user preferences to suggest relevant items from large catalogs.\n\nTopics: collaborative filtering, matrix factorization, content-based methods, hybrid approaches.",
    "Deploying ML models requires infrastructure for serving predictions, monitoring, and maintenance.\n\nTopics: REST APIs, containerization, A/B testing, model drift detection, CI/CD for ML.",
    "Current research explores the boundaries of what machine learning can achieve across diverse domains.\n\nTopics: foundation models, multimodal learning, federated learning, neural architecture search.",
    "This course provided a comprehensive overview of machine learning theory and practice.\n\nKey takeaways: algorithm selection, evaluation methodology, ethical considerations, practical deployment.",
]


def set_master_body_style(prs, font_name='Liberation Sans', font_size_pt=18):
    """Set the slide master body text placeholder style via XML."""
    slide_master = prs.slide_masters[0]
    # Access the txStyles element in the slide master
    master_el = slide_master.element
    nsmap = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    # Find or create txStyles
    txStyles = master_el.find('.//p:txStyles', nsmap)
    if txStyles is None:
        txStyles = master_el.makeelement(qn('p:txStyles'), {})
        master_el.append(txStyles)

    # Find or create bodyStyle
    bodyStyle = txStyles.find('p:bodyStyle', nsmap)
    if bodyStyle is None:
        bodyStyle = txStyles.makeelement(qn('p:bodyStyle'), {})
        txStyles.append(bodyStyle)

    # Clear existing lvl1pPr if any, then create new one
    for existing in bodyStyle.findall('a:lvl1pPr', nsmap):
        bodyStyle.remove(existing)

    lvl1pPr = bodyStyle.makeelement(qn('a:lvl1pPr'), {})
    bodyStyle.append(lvl1pPr)

    # Set font
    defRPr = lvl1pPr.makeelement(qn('a:defRPr'), {
        'sz': str(font_size_pt * 100),  # size in hundredths of a point
        'lang': 'en-US',
    })
    lvl1pPr.append(defRPr)

    # Set font name via latin element
    latin = defRPr.makeelement(qn('a:latin'), {'typeface': font_name})
    defRPr.append(latin)

    # Single line spacing (100%) - this is the default, set explicitly
    lnSpc = lvl1pPr.makeelement(qn('a:lnSpc'), {})
    spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': '100000'})  # 100% = single
    lnSpc.append(spcPct)
    lvl1pPr.append(lnSpc)

    # Also set body placeholder shapes on the master to use Liberation Sans 18pt
    for shape in slide_master.shapes:
        if shape.has_text_frame:
            ph = shape.placeholder_format
            if ph is not None and ph.idx == 1:  # body placeholder
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = font_name
                        run.font.size = Pt(font_size_pt)


def create_initial():
    prs = Presentation()

    # Set master body style to Liberation Sans 18pt, single spacing
    set_master_body_style(prs, 'Liberation Sans', 18)

    # Also set it on all slide layouts' body placeholders
    for layout in prs.slide_layouts:
        for shape in layout.placeholders:
            if shape.placeholder_format.idx == 1:  # body placeholder
                for para in shape.text_frame.paragraphs:
                    pPr = para._p.get_or_add_pPr()
                    # Set font
                    defRPr = pPr.find(qn('a:defRPr'))
                    if defRPr is None:
                        defRPr = pPr.makeelement(qn('a:defRPr'), {
                            'sz': '1800',
                            'lang': 'en-US',
                        })
                        pPr.append(defRPr)
                    else:
                        defRPr.set('sz', '1800')
                    # Set latin font
                    latin = defRPr.find(qn('a:latin'))
                    if latin is None:
                        latin = defRPr.makeelement(qn('a:latin'), {'typeface': 'Liberation Sans'})
                        defRPr.append(latin)
                    else:
                        latin.set('typeface', 'Liberation Sans')

    # Create 30 slides
    for i in range(30):
        title_text, subtitle = LECTURE_TOPICS[i]
        body_text = BODY_CONTENTS[i]

        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle
        else:
            # Title + Content layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Lecture {i}: {title_text}"
            if len(slide.placeholders) > 1:
                body_ph = slide.placeholders[1]
                tf = body_ph.text_frame
                tf.clear()
                # Set body text with Liberation Sans 18pt explicitly
                paragraphs = body_text.split('\n\n')
                for pi, para_text in enumerate(paragraphs):
                    if pi == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    run = p.add_run()
                    run.text = para_text
                    run.font.name = 'Liberation Sans'
                    run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
