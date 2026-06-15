"""
Initial Setup: Talent show presentation with Fade transition on slide 5 (no sound)
Task ID: impress_tm_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16):
    """Add multiple bullet points as paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)


def set_slide_background(slide, r, g, b):
    """Set solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_fade_transition(slide, duration_ms=2000):
    """Add a Fade transition to a slide via XML manipulation."""
    slide_elem = slide._element
    # Remove any existing transition
    for existing in slide_elem.findall(qn('p:transition')):
        slide_elem.remove(existing)
    # Add transition element
    transition = etree.SubElement(slide_elem, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advTm', str(duration_ms))
    fade = etree.SubElement(transition, qn('p:fade'))


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Riverside Community Talent Show 2025"
    slide1.placeholders[1].text = "Saturday, June 14th | Riverside Auditorium"
    set_slide_background(slide1, 0x1B, 0x2A, 0x4A)
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 2: Schedule Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, 0.5, 0.3, 9, 1, "Event Schedule", font_size=32,
                bold=True, color=(0x1B, 0x2A, 0x4A), alignment=PP_ALIGN.CENTER)
    schedule_items = [
        "5:00 PM - Doors Open & Refreshments",
        "5:30 PM - Welcome Address by Mayor Linda Torres",
        "6:00 PM - Act 1: Musical Performances",
        "6:45 PM - Intermission (15 min)",
        "7:00 PM - Act 2: Dance & Comedy",
        "7:45 PM - Act 3: Special Performances",
        "8:15 PM - Judges' Deliberation",
        "8:30 PM - Awards Ceremony & Closing",
    ]
    add_bullet_points(slide2, 1.0, 1.5, 8, 5.5, schedule_items, font_size=16)

    # --- Slide 3: Judges Panel ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, 0.5, 0.3, 9, 1, "Meet Our Judges", font_size=32,
                bold=True, color=(0x1B, 0x2A, 0x4A), alignment=PP_ALIGN.CENTER)
    judges = [
        "Dr. Patricia Owens - Music Department Chair, Riverside University",
        "Marcus DeLeon - Broadway choreographer & dance instructor",
        "Jennifer Nakamura - Comedy writer, The Laughing Room",
        "David Okonkwo - Local radio host, WRIV 98.5 FM",
    ]
    add_bullet_points(slide3, 0.8, 1.5, 8.4, 4, judges, font_size=16)
    add_textbox(slide3, 0.8, 5.5, 8.4, 1.5,
                "Judging Criteria: Originality (30%), Skill (30%), Stage Presence (20%), Audience Connection (20%)",
                font_size=14, color=(0x66, 0x66, 0x66))

    # --- Slide 4: Act 1 - Musical Performances ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide4, 0xF0, 0xF4, 0xF8)
    add_textbox(slide4, 0.5, 0.3, 9, 1, "Act 1: Musical Performances", font_size=28,
                bold=True, color=(0x2E, 0x7D, 0x32))
    act1_performers = [
        "1. Emily Vasquez - Classical violin solo (Paganini Caprice No. 24)",
        "2. The Riverside Trio - Jazz ensemble (Original composition: 'Evening Glow')",
        "3. Aiden Park - Acoustic guitar & vocals ('Riverside Sunset')",
        "4. Sofia Ramirez - Piano recital (Chopin Ballade No. 1)",
        "5. The Harmonic Five - A cappella group ('Bohemian Rhapsody' arrangement)",
    ]
    add_bullet_points(slide4, 0.8, 1.5, 8.4, 5, act1_performers, font_size=15)

    # --- Slide 5: Act 2 - Dance & Comedy (THIS SLIDE GETS FADE TRANSITION) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide5, 0xFD, 0xF0, 0xE6)
    add_textbox(slide5, 0.5, 0.3, 9, 1, "Act 2: Dance & Comedy", font_size=28,
                bold=True, color=(0xE6, 0x51, 0x00))
    act2_performers = [
        "6. Luna Chen & Jake Morrison - Contemporary dance duet ('Gravity')",
        "7. Marcus Williams - Stand-up comedy (10 min set)",
        "8. The Breakaway Crew - Hip-hop dance group (6 members)",
        "9. Rachel Okafor - Improv comedy with audience participation",
        "10. Priya Sharma & Nadia Kozlov - Ballroom dance (Viennese Waltz)",
    ]
    add_bullet_points(slide5, 0.8, 1.5, 8.4, 5, act2_performers, font_size=15)
    # Add Fade transition to slide 5 with 2.0s duration, NO sound
    add_fade_transition(slide5, duration_ms=2000)

    # --- Slide 6: Act 3 - Special Performances ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide6, 0xE8, 0xEA, 0xF6)
    add_textbox(slide6, 0.5, 0.3, 9, 1, "Act 3: Special Performances", font_size=28,
                bold=True, color=(0x4A, 0x14, 0x8C))
    act3_performers = [
        "11. The Riverside Magic Circle - Stage illusions (15 min)",
        "12. Tomoko Hayashi - Traditional Japanese taiko drumming",
        "13. Alex Rivera - Spoken word poetry ('Voices of Our Town')",
        "14. The Flying Sparks - Circus arts (aerial silks & juggling)",
    ]
    add_bullet_points(slide6, 0.8, 1.5, 8.4, 5, act3_performers, font_size=15)

    # --- Slide 7: Prizes & Sponsors ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, 0.5, 0.3, 9, 1, "Prizes & Sponsors", font_size=32,
                bold=True, color=(0x1B, 0x2A, 0x4A), alignment=PP_ALIGN.CENTER)
    prizes = [
        "Grand Prize: $2,500 + Recording session at Riverside Studios",
        "Runner-Up: $1,000 + Feature in Riverside Weekly",
        "People's Choice: $500 + Gift basket from local businesses",
        "Best New Talent (Under 18): $750 scholarship",
    ]
    add_bullet_points(slide7, 0.8, 1.5, 8.4, 3, prizes, font_size=16)
    add_textbox(slide7, 0.8, 5.0, 8.4, 1.5,
                "Proudly sponsored by: Riverside Credit Union | Torres Family Foundation | WRIV Radio | Riverside Studios",
                font_size=13, color=(0x88, 0x88, 0x88))

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_background(slide8, 0x1B, 0x2A, 0x4A)
    add_textbox(slide8, 1, 2.5, 8, 2, "Thank You for Attending!", font_size=36,
                bold=True, color=(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide8, 1.5, 4.5, 7, 1.5,
                "Follow us @RiversideTalentShow | riversidetalentshow.org",
                font_size=18, color=(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
