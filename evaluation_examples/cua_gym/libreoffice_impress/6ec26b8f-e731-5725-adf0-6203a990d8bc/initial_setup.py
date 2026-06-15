"""
Initial Setup: Create a 4-slide presentation with title and 3 bullet text boxes on slide 2, no animations.
Task ID: impress_ma_082
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_082'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_name="Calibri", font_size=Pt(18),
                      bold=False, color=None, alignment=None):
    """Helper to set text properties on a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    # Standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ========== SLIDE 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # Title
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Coordinated Marketing Strategy"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(2.5), Inches(3.8), Inches(8), Inches(1.0))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Q3 2025 Campaign Planning & Resource Allocation"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)

    # Date
    date_box = slide1.shapes.add_textbox(Inches(4.5), Inches(5.2), Inches(4), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Presented by Elena Rodriguez  |  July 15, 2025"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x99, 0xAA, 0xBB)

    # ========== SLIDE 2: Key Objectives (Title + 3 bullet text boxes) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide2.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # Title for slide 2
    title2 = slide2.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11), Inches(1.0))
    tf = title2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Campaign Objectives"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # Bullet 1 text box
    bullet1 = slide2.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.2))
    tf = bullet1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Increase brand awareness by 35% across digital channels through targeted social media campaigns and influencer partnerships"
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Bullet 2 text box
    bullet2 = slide2.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10), Inches(1.2))
    tf = bullet2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Drive 50,000 qualified leads to the product landing page with an expected conversion rate of 4.2% by end of Q3"
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Bullet 3 text box
    bullet3 = slide2.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10), Inches(1.2))
    tf = bullet3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Establish strategic co-marketing agreements with at least 3 enterprise partners to expand market reach into the APAC region"
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ========== SLIDE 3: Budget Overview ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    title3 = slide3.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11), Inches(1.0))
    tf = title3.text_frame
    p = tf.paragraphs[0]
    p.text = "Budget Allocation Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # Add table
    table_shape = slide3.shapes.add_table(5, 3, Inches(1.5), Inches(2.0), Inches(10), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)

    headers = ["Category", "Allocated Budget", "% of Total"]
    data = [
        ["Digital Advertising", "$125,000", "41.7%"],
        ["Content Production", "$85,000", "28.3%"],
        ["Events & Sponsorships", "$55,000", "18.3%"],
        ["Analytics & Tools", "$35,000", "11.7%"],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)

    # ========== SLIDE 4: Timeline ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    title4 = slide4.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11), Inches(1.0))
    tf = title4.text_frame
    p = tf.paragraphs[0]
    p.text = "Implementation Timeline"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Calibri"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    # Timeline items as text boxes
    phases = [
        ("Phase 1: Research & Planning", "July 1 - July 21, 2025", "Market analysis, competitor benchmarking, audience segmentation"),
        ("Phase 2: Content Development", "July 22 - August 15, 2025", "Creative asset production, landing page design, A/B test variants"),
        ("Phase 3: Campaign Launch", "August 18 - September 15, 2025", "Multi-channel campaign rollout, real-time performance monitoring"),
        ("Phase 4: Analysis & Reporting", "September 16 - September 30, 2025", "ROI analysis, stakeholder presentations, Q4 recommendations"),
    ]

    for i, (phase_title, dates, desc) in enumerate(phases):
        y_pos = 1.8 + i * 1.3
        # Phase title
        box = slide4.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = phase_title
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

        # Dates + description
        desc_box = slide4.shapes.add_textbox(Inches(1.5), Inches(y_pos + 0.4), Inches(10), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{dates} -- {desc}"
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
