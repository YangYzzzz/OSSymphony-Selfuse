"""
Initial Setup: Create Product Review presentation with 8 slides.
Task ID: impress_exec_065
Domain: libreoffice_impress
Slide 4 has title 'Product Line Performance' but NO chart.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_065'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, layout_idx, title_text, body_lines):
    """Helper to add a slide with title and bullet body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    if body_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Product Review 2025"
    slide1.placeholders[1].text = "Prepared by Strategic Planning Division\nQ4 Executive Summary"

    # ---- Slide 2: Market Overview ----
    add_title_body_slide(prs, 1, "Market Overview", [
        "Global SaaS market grew 18% YoY to $232B",
        "North America accounts for 52% of total addressable market",
        "APAC region showing fastest growth at 24% CAGR",
        "Key drivers: cloud adoption, remote work, AI integration",
        "Competitive intensity increased with 14 new entrants",
    ])

    # ---- Slide 3: Competitive Landscape ----
    add_title_body_slide(prs, 1, "Competitive Landscape", [
        "Market share: Our position improved from #4 to #3",
        "Leader: TechCorp (28%) - strong enterprise foothold",
        "Runner-up: CloudFirst (22%) - aggressive pricing",
        "Our share: 15% (+3pp YoY) driven by SMB expansion",
        "Key differentiator: integrated analytics platform",
        "Threat: CloudFirst acquiring DataViz Inc.",
    ])

    # ---- Slide 4: Product Line Performance (NO CHART) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Line Performance"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Subtitle text below the title
    txBox2 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Quarterly revenue breakdown by product line will be visualized here"
    run2 = p2.runs[0]
    run2.font.size = Pt(14)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ---- Slide 5: Customer Satisfaction ----
    add_title_body_slide(prs, 1, "Customer Satisfaction", [
        "Overall NPS improved from 42 to 58 (+16 points)",
        "Enterprise segment NPS: 65 (industry-leading)",
        "SMB segment NPS: 54 (up from 38)",
        "Consumer segment NPS: 48 (needs attention)",
        "Top feedback: onboarding speed, API documentation",
        "Action: Launched Customer Success Academy in Q3",
    ])

    # ---- Slide 6: Growth Strategy ----
    add_title_body_slide(prs, 1, "Growth Strategy 2026", [
        "Pillar 1: Enterprise upsell — target $15M ARR from existing accounts",
        "Pillar 2: SMB self-serve — reduce CAC by 30% with product-led growth",
        "Pillar 3: International expansion — EMEA launch Q2 2026",
        "Pillar 4: AI-powered analytics — beta in Q1, GA in Q3",
        "Investment: $12M allocated across R&D and GTM",
    ])

    # ---- Slide 7: Financial Projections ----
    add_title_body_slide(prs, 1, "Financial Projections", [
        "2025 Total Revenue: $41.5M (Enterprise $41.5M combined)",
        "2026 Target: $58M (+40% growth)",
        "Gross margin target: 78% (up from 74%)",
        "Path to profitability: Q3 2026 breakeven",
        "Key assumption: 90% enterprise renewal rate maintained",
    ])

    # ---- Slide 8: Next Steps ----
    add_title_body_slide(prs, 1, "Next Steps & Action Items", [
        "Complete Q1 2026 product roadmap by January 15",
        "Finalize EMEA go-to-market strategy by February 1",
        "Board presentation scheduled for January 28",
        "Hiring plan: 35 new roles across Engineering and Sales",
        "Monthly review cadence starting February",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
