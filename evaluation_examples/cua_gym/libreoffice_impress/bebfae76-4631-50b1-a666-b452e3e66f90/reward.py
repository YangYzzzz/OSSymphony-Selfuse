"""
Reward Script: Add multi-series line chart on slide 4 comparing revenue across product lines
Task ID: impress_exec_065
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Chart exists on slide 4
  Component 2 (0.15): Chart is LINE type with 3 series
  Component 3 (0.15): Chart title is 'Revenue by Product Line'
  Component 4 (0.10): Series names match (Enterprise, SMB, Consumer)
  Component 5 (0.20): Data values correct for all 3 series
  Component 6 (0.20): Series colors correct (#003366, #2196F3, #90CAF9)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_065'

EXPECTED_SERIES = {
    'Enterprise': {
        'values': [8.0, 9.5, 11.0, 13.0],
        'color': '003366',
    },
    'SMB': {
        'values': [3.0, 3.5, 3.8, 4.2],
        'color': '2196F3',
    },
    'Consumer': {
        'values': [1.0, 1.5, 1.4, 1.6],
        'color': '90CAF9',
    },
}


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice edits."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 4 (0.20 points)
    try:
        if chart_shape is not None:
            print(f"PASS: Component 1 — Chart found on slide 4 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — No chart found on slide 4")
            # No chart means nothing else can pass
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart is LINE type with 3 series (0.15 points)
    try:
        chart_type = chart.chart_type
        num_series = len(chart.series)
        # LINE types: LINE (65=LINE_MARKERS, 64=LINE, etc.)
        # Accept any line-type chart
        is_line = chart_type is not None and 'LINE' in str(chart_type).upper()
        if is_line and num_series == 3:
            print(f"PASS: Component 2 — Line chart with 3 series (type={chart_type}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Expected LINE chart with 3 series, got type={chart_type}, series={num_series}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart title is 'Revenue by Product Line' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Revenue by Product Line':
                print(f"PASS: Component 3 — Chart title matches exactly (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 — Expected title 'Revenue by Product Line', found '{title_text}'")
        else:
            print(f"FAIL: Component 3 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Extract series info via XML for remaining components
    ns = {
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }

    series_data = {}
    try:
        for i, series in enumerate(chart.series):
            ser_el = series._element
            tx = ser_el.find('.//c:tx//c:v', ns)
            name = tx.text.strip() if tx is not None else f"Series_{i}"

            data_vals = ser_el.findall('.//c:val//c:v', ns)
            values = [float(v.text) for v in data_vals]

            a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            all_srgb = ser_el.findall(f'.//{{{a_ns}}}srgbClr')
            colors = set()
            for clr in all_srgb:
                color_val = clr.get('val')
                if color_val:
                    colors.add(color_val.upper())

            series_data[name] = {'values': values, 'colors': colors}
            print(f"  Found series '{name}': values={values}, colors={colors}")
    except Exception as e:
        print(f"ERROR: Extracting series data — {e}")

    # Component 4: Series names match (0.10 points)
    try:
        expected_names = set(EXPECTED_SERIES.keys())
        actual_names = set(series_data.keys())
        if expected_names == actual_names:
            print(f"PASS: Component 4 — Series names match: {actual_names} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — Expected names {expected_names}, found {actual_names}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Data values correct for all 3 series (0.20 points)
    # Award partial: ~0.067 per correct series
    try:
        correct_series_count = 0
        for name, expected in EXPECTED_SERIES.items():
            if name in series_data:
                actual_vals = series_data[name]['values']
                expected_vals = expected['values']
                if len(actual_vals) == len(expected_vals):
                    all_match = all(
                        abs(a - e) < 0.01 for a, e in zip(actual_vals, expected_vals)
                    )
                    if all_match:
                        correct_series_count += 1
                        print(f"  PASS: {name} values match: {actual_vals}")
                    else:
                        print(f"  FAIL: {name} values mismatch: expected {expected_vals}, got {actual_vals}")
                else:
                    print(f"  FAIL: {name} value count mismatch: expected {len(expected_vals)}, got {len(actual_vals)}")
            else:
                print(f"  FAIL: {name} series not found")

        if correct_series_count == 3:
            print(f"PASS: Component 5 — All 3 series data values correct (0.20 pts)")
            total_score += 0.20
        elif correct_series_count > 0:
            partial = round(0.20 * correct_series_count / 3, 2)
            print(f"PARTIAL: Component 5 — {correct_series_count}/3 series correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No series data values matched")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Series colors correct (0.20 points)
    # Award partial: ~0.067 per correct series color
    try:
        correct_color_count = 0
        for name, expected in EXPECTED_SERIES.items():
            expected_color = expected['color'].upper()
            if name in series_data:
                actual_colors = series_data[name]['colors']
                if expected_color in actual_colors:
                    correct_color_count += 1
                    print(f"  PASS: {name} color {expected_color} found")
                else:
                    print(f"  FAIL: {name} expected color {expected_color}, found {actual_colors}")
            else:
                print(f"  FAIL: {name} series not found for color check")

        if correct_color_count == 3:
            print(f"PASS: Component 6 — All 3 series colors correct (0.20 pts)")
            total_score += 0.20
        elif correct_color_count > 0:
            partial = round(0.20 * correct_color_count / 3, 2)
            print(f"PARTIAL: Component 6 — {correct_color_count}/3 colors correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 6 — No series colors matched")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
