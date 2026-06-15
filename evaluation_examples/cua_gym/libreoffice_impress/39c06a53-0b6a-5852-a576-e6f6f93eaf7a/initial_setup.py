"""
Initial Setup: Brand presentation with white background and black title text on slide master.
Task ID: impress_fix_074
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_074'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # -------------------------------------------------------------------
    # Set the slide master background to WHITE (explicit solid fill)
    # -------------------------------------------------------------------
    master = prs.slide_masters[0]
    # master.background._element is <p:cSld>; navigate to <p:bg> child
    cSld = master.background._element
    bg = cSld.find(qn('p:bg'))
    if bg is None:
        bg = cSld.makeelement(qn('p:bg'), {})
        cSld.insert(0, bg)
    # Remove bgRef (theme reference overrides explicit fill)
    bgRef = bg.find(qn('p:bgRef'))
    if bgRef is not None:
        bg.remove(bgRef)
    # Remove existing bgPr if any
    bgPr = bg.find(qn('p:bgPr'))
    if bgPr is not None:
        bg.remove(bgPr)
    # Add new bgPr with solid white fill
    bgPr = bg.makeelement(qn('p:bgPr'), {})
    solidFill = bgPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': 'FFFFFF'})
    solidFill.append(srgbClr)
    bgPr.append(solidFill)
    bg.append(bgPr)

    # -------------------------------------------------------------------
    # Set the slide master title placeholder text color to BLACK (#000000)
    # -------------------------------------------------------------------
    for sp in master.placeholders:
        if sp.placeholder_format.idx == 0:  # Title placeholder
            # Set default run properties on the placeholder's text body
            txBody = sp._element.find(qn('p:txBody'))
            if txBody is None:
                continue
            lstStyle = txBody.find(qn('a:lstStyle'))
            if lstStyle is None:
                lstStyle = txBody.makeelement(qn('a:lstStyle'), {})
                txBody.insert(0, lstStyle)
            # Set level 0 default run properties
            lvl1pPr = lstStyle.find(qn('a:lvl1pPr'))
            if lvl1pPr is None:
                lvl1pPr = lstStyle.makeelement(qn('a:lvl1pPr'), {})
                lstStyle.append(lvl1pPr)
            defRPr = lvl1pPr.find(qn('a:defRPr'))
            if defRPr is None:
                defRPr = lvl1pPr.makeelement(qn('a:defRPr'), {'sz': '4400'})
                lvl1pPr.append(defRPr)
            # Remove any existing fill in defRPr
            for sf in defRPr.findall(qn('a:solidFill')):
                defRPr.remove(sf)
            # Add black color
            sf = defRPr.makeelement(qn('a:solidFill'), {})
            clr = sf.makeelement(qn('a:srgbClr'), {'val': '000000'})
            sf.append(clr)
            defRPr.insert(0, sf)
            break

    # -------------------------------------------------------------------
    # Slide content data - realistic brand/rebrand presentation
    # -------------------------------------------------------------------
    slide_content = [
        {
            'layout': 0,
            'title': 'Nexora Global Rebrand Initiative',
            'subtitle': 'Strategic Brand Refresh for FY2026\nPrepared by the Marketing Strategy Team'
        },
        {
            'layout': 1,
            'title': 'Executive Summary',
            'body': 'Nexora Global is undergoing a comprehensive brand refresh to better align with our evolved mission and market positioning. This presentation outlines the strategic rationale, visual identity changes, and implementation timeline for the rebrand across all customer touchpoints.'
        },
        {
            'layout': 1,
            'title': 'Brand Evolution Timeline',
            'body': '2018: Company founded as NexaTech Solutions\n2020: Expanded to enterprise markets, renamed Nexora\n2022: Acquired CloudPulse and DataStream\n2024: Global expansion into APAC and EMEA\n2026: Brand refresh to unify acquired entities'
        },
        {
            'layout': 1,
            'title': 'Market Research Findings',
            'body': '78% of customers associate Nexora with innovation\n62% find current branding inconsistent across products\n91% prefer modern, clean visual identity\nKey competitor analysis shows trend toward minimalist design'
        },
        {
            'layout': 1,
            'title': 'Core Brand Values',
            'body': 'Innovation: Pioneering solutions that transform industries\nTrust: Building lasting relationships through transparency\nExcellence: Delivering quality in every interaction\nCollaboration: Growing stronger through partnership'
        },
        {
            'layout': 1,
            'title': 'Visual Identity Overview',
            'body': 'Primary color palette update\nTypography system refinement\nLogo mark simplification\nIconography style guide\nPhotography direction'
        },
        {
            'layout': 1,
            'title': 'Color Palette Strategy',
            'body': 'Primary: Deep Navy and Silver\nSecondary: Warm Teal and Soft Gold\nAccent: Coral for CTAs and highlights\nNeutral: Slate Gray gradient system\nAccessibility: All combinations meet WCAG 2.1 AA'
        },
        {
            'layout': 1,
            'title': 'Typography Guidelines',
            'body': 'Headlines: Montserrat Bold (24-44pt)\nSubheadings: Montserrat SemiBold (18-22pt)\nBody text: Open Sans Regular (11-14pt)\nCode/Data: JetBrains Mono (10-12pt)\nAll fonts available via Google Fonts for web consistency'
        },
        {
            'layout': 1,
            'title': 'Logo Redesign Rationale',
            'body': 'Simplified geometric mark for digital-first applications\nScalable from 16px favicon to 4m billboard\nMonochrome variant for single-color contexts\nAnimated version for digital platforms\nRetains brand equity from previous hexagonal motif'
        },
        {
            'layout': 1,
            'title': 'Digital Touchpoint Updates',
            'body': 'Corporate website: Full redesign launching Q2 2026\nMobile apps: Phased icon and UI updates\nEmail templates: New header and footer system\nSocial media: Updated profile assets and templates\nProduct dashboards: Gradual theme migration'
        },
        {
            'layout': 1,
            'title': 'Physical Collateral Changes',
            'body': 'Business cards: New layout with QR code integration\nLetterhead and envelopes: Updated header design\nOffice signage: Phased replacement in 23 locations\nTrade show booth: Redesigned modular system\nMerchandise: Refreshed product catalog'
        },
        {
            'layout': 1,
            'title': 'Implementation Roadmap',
            'body': 'Phase 1 (Q1 2026): Internal launch and employee training\nPhase 2 (Q2 2026): Digital properties migration\nPhase 3 (Q3 2026): Customer communications\nPhase 4 (Q4 2026): Physical assets and signage\nPhase 5 (Q1 2027): Full audit and compliance check'
        },
        {
            'layout': 1,
            'title': 'Budget Allocation',
            'body': 'Design and Creative: $340,000 (22%)\nDigital Implementation: $520,000 (34%)\nPhysical Collateral: $280,000 (18%)\nTraining and Change Management: $150,000 (10%)\nContingency and QA: $245,000 (16%)\nTotal Budget: $1,535,000'
        },
        {
            'layout': 1,
            'title': 'Risk Assessment',
            'body': 'Brand confusion during transition (Medium risk)\nCustomer resistance to change (Low risk)\nInternal adoption delays (Medium risk)\nVendor timeline dependencies (High risk)\nBudget overrun on physical assets (Low risk)'
        },
        {
            'layout': 1,
            'title': 'Success Metrics',
            'body': 'Brand recognition increase: Target +15% by Q4 2026\nCustomer satisfaction (CSAT): Maintain >4.2/5.0\nEmployee brand ambassador score: Target >85%\nWeb traffic from brand searches: +25%\nSocial media engagement rate: +30%'
        },
        {
            'layout': 1,
            'title': 'Stakeholder Communication Plan',
            'body': 'Board of Directors: Monthly progress briefings\nSenior Leadership: Bi-weekly status updates\nAll Employees: Town hall + Intranet portal\nKey Customers: Personal outreach from account managers\nMedia: Press release timed with public launch'
        },
        {
            'layout': 1,
            'title': 'Competitive Landscape',
            'body': 'Meridian Corp: Rebranded 2024, 12% market share gain\nVelocity Systems: Ongoing refresh, mixed reception\nAtlas Dynamics: Strong legacy brand, no recent changes\nPinnacle Tech: Minimalist rebrand 2023, positive sentiment\nNexora must differentiate through cohesive narrative'
        },
        {
            'layout': 1,
            'title': 'Agency Partners',
            'body': 'Lead Creative: Sterling & Associates (NYC)\nDigital Strategy: PixelWave Interactive (SF)\nPR and Comms: Meridian Public Relations (London)\nPrint Production: Apex Printing Solutions (Chicago)\nPhotography: Lumen Studios (LA)'
        },
        {
            'layout': 1,
            'title': 'Key Milestones',
            'body': 'Feb 15: Final brand guide approval\nMar 1: Employee launch event\nApr 15: Website beta with new branding\nJun 1: Public launch across all digital channels\nSep 30: Physical signage 80% complete\nDec 31: Full compliance audit'
        },
        {
            'layout': 1,
            'title': 'Next Steps and Action Items',
            'body': 'Finalize color palette with stakeholder sign-off\nComplete typography licensing agreements\nBegin website redesign sprint (2-week cycles)\nSchedule employee brand training sessions\nCoordinate vendor timelines for Q3 physical rollout\nEstablish brand compliance monitoring dashboard'
        },
    ]

    # -------------------------------------------------------------------
    # Create all 20 slides
    # -------------------------------------------------------------------
    for i, content in enumerate(slide_content):
        layout_idx = content['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = content['title']

        # Body or subtitle
        if layout_idx == 0 and 'subtitle' in content:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content['subtitle']
        elif 'body' in content:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content['body']

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Number of slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
