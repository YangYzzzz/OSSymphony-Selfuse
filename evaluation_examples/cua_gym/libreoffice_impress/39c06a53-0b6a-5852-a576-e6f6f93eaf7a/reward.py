"""
Reward Script: Change slide master background to #F5F5F5 and title text color to #003366
Task ID: impress_fix_074
Domain: libreoffice_impress
Scoring:
  Component 1: Master background is #F5F5F5 (0.4 pts)
  Component 2: Master title placeholder text color is #003366 (0.4 pts)
  Component 3: All 20 slides present and inherit master bg (0.2 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_074'

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_master_bg_color(pptx_path):
    """Extract master slide background srgbClr from XML."""
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
            root = ET.parse(f).getroot()
        bg = root.find('.//p:cSld/p:bg', NS)
        if bg is not None:
            for clr in bg.findall('.//a:srgbClr', NS):
                return clr.get('val')
    return None


def get_master_title_color(pptx_path):
    """Extract title placeholder defRPr srgbClr from slide master XML."""
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
            root = ET.parse(f).getroot()
        for sp in root.findall('.//p:cSld/p:spTree/p:sp', NS):
            ph = sp.find('.//p:nvSpPr/p:nvPr/p:ph', NS)
            if ph is not None and ph.get('type') == 'title':
                # Check defRPr solidFill srgbClr (the default run property for titles)
                for defRPr in sp.findall('.//a:defRPr', NS):
                    for sc in defRPr.findall('.//a:srgbClr', NS):
                        return sc.get('val')
                # Fallback: check any srgbClr in the title shape
                for sc in sp.findall('.//a:srgbClr', NS):
                    return sc.get('val')
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be a valid pptx
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Master slide background is #F5F5F5 (0.4 points)
    try:
        bg_color = get_master_bg_color(file_path)
        if bg_color is not None and bg_color.upper() == 'F5F5F5':
            print(f"PASS: Component 1 -- Master BG color is F5F5F5 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- Expected master BG F5F5F5, found: {bg_color}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Master title placeholder text color is #003366 (0.4 points)
    try:
        title_color = get_master_title_color(file_path)
        if title_color is not None and title_color.upper() == '003366':
            print(f"PASS: Component 2 -- Master title color is 003366 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 2 -- Expected title color 003366, found: {title_color}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: All 20 slides present and inheriting master bg (0.2 points)
    # This verifies the changes propagate - slides must still exist and
    # should NOT have overridden backgrounds (type 5 = inherited from master)
    try:
        slide_count = len(prs.slides)
        if slide_count != 20:
            print(f"FAIL: Component 3 -- Expected 20 slides, found {slide_count}")
        else:
            # Check that slides inherit from master (bg fill type 5)
            # AND that the inherited master bg is F5F5F5 (not FFFFFF)
            # This component only awards points if the bg was actually changed
            # (i.e., it depends on Component 1 passing too)
            master_bg = get_master_bg_color(file_path)
            if master_bg is not None and master_bg.upper() == 'F5F5F5':
                # Verify at least a sample of slides inherit from master
                inherited_count = 0
                for i, slide in enumerate(prs.slides):
                    sfill = slide.background.fill
                    # fill type 5 = BACKGROUND (inherited) or None
                    if sfill.type is None or sfill.type == 5:
                        inherited_count += 1
                if inherited_count == 20:
                    print(f"PASS: Component 3 -- All 20 slides present, all inherit master BG F5F5F5 (0.2 pts)")
                    total_score += 0.2
                else:
                    print(f"FAIL: Component 3 -- {inherited_count}/20 slides inherit master BG; {20 - inherited_count} have overridden BG")
            else:
                print(f"FAIL: Component 3 -- Master BG is not F5F5F5, so inherited slides don't reflect change")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
