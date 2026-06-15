"""
Reward Script: Set content text boxes on slide 4 to left-aligned, title to center-aligned.
Task ID: impress_teach_016
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): TextBox 3 (content) all paragraphs LEFT-aligned (was CENTER)
  Component 2 (0.35): TextBox 4 (content) all paragraphs LEFT-aligned (was CENTER)
  Component 3 (0.15): Title text box (TextBox 2) still CENTER-aligned AND content boxes are LEFT
                       (compound check anchored to the change)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_016'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Identify shapes by name and type
    # Shape 0: "Title 1" (placeholder) - empty title
    # Shape 1: "TextBox 2" - visual title "Memory Systems and Encoding" (should be CENTER)
    # Shape 2: "TextBox 3" - content (should be LEFT after task)
    # Shape 3: "TextBox 4" - content (should be LEFT after task)

    # Build a map of text boxes (non-placeholder shapes with text)
    textboxes = []
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            # Check if it's a placeholder title
            try:
                pf = shape.placeholder_format
                if pf.type is not None and pf.idx == 0:
                    # This is the placeholder title (empty in this presentation)
                    continue
            except (ValueError, AttributeError):
                pass
            textboxes.append(shape)

    if len(textboxes) < 3:
        print(f"CRITICAL: Expected at least 3 text boxes on slide 4, found {len(textboxes)}")
        print("REWARD: 0.0")
        return 0.0

    # The first textbox is the visual title, remaining are content
    title_tb = textboxes[0]
    content_tbs = textboxes[1:]

    def normalize_alignment(align):
        """Normalize None to LEFT (default alignment)."""
        if align is None:
            return PP_ALIGN.LEFT
        return align

    # Component 1: TextBox 3 (first content box) - all paragraphs LEFT-aligned (0.5 points)
    # This was CENTER in initial, should be LEFT in golden.
    try:
        content_box_1 = content_tbs[0]
        paras = [p for p in content_box_1.text_frame.paragraphs if (p.text or "").strip()]
        if len(paras) == 0:
            print(f"FAIL: Component 1 — No non-empty paragraphs in first content textbox")
        else:
            all_left = True
            for p in paras:
                actual = normalize_alignment(p.alignment)
                if actual != PP_ALIGN.LEFT:
                    all_left = False
                    print(f"FAIL: Component 1 — Paragraph '{p.text[:40]}' has alignment {p.alignment}, expected LEFT")
                    break
            if all_left:
                print(f"PASS: Component 1 — All {len(paras)} paragraphs in first content box are LEFT-aligned (0.5 pts)")
                total_score += 0.5
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: TextBox 4 (second content box) - all paragraphs LEFT-aligned (0.35 points)
    # This was CENTER in initial, should be LEFT in golden.
    try:
        content_box_2 = content_tbs[1]
        paras = [p for p in content_box_2.text_frame.paragraphs if (p.text or "").strip()]
        if len(paras) == 0:
            print(f"FAIL: Component 2 — No non-empty paragraphs in second content textbox")
        else:
            all_left = True
            for p in paras:
                actual = normalize_alignment(p.alignment)
                if actual != PP_ALIGN.LEFT:
                    all_left = False
                    print(f"FAIL: Component 2 — Paragraph '{p.text[:40]}' has alignment {p.alignment}, expected LEFT")
                    break
            if all_left:
                print(f"PASS: Component 2 — All {len(paras)} paragraphs in second content box are LEFT-aligned (0.35 pts)")
                total_score += 0.35
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Title textbox stays CENTER AND at least one content box is LEFT (0.15 points)
    # Compound check: the title being CENTER only earns points if content was actually changed to LEFT.
    # This prevents scoring on initial_env where everything is CENTER.
    try:
        title_paras = [p for p in title_tb.text_frame.paragraphs if (p.text or "").strip()]
        title_centered = True
        for p in title_paras:
            actual = normalize_alignment(p.alignment)
            if actual != PP_ALIGN.CENTER:
                title_centered = False
                print(f"FAIL: Component 3 — Title paragraph '{p.text[:40]}' has alignment {p.alignment}, expected CENTER")
                break

        # Check that at least one content box is LEFT (anchoring to the change)
        any_content_left = False
        for cb in content_tbs:
            cb_paras = [p for p in cb.text_frame.paragraphs if (p.text or "").strip()]
            if cb_paras and all(normalize_alignment(p.alignment) == PP_ALIGN.LEFT for p in cb_paras):
                any_content_left = True
                break

        if title_centered and any_content_left:
            print(f"PASS: Component 3 — Title is CENTER-aligned and content is LEFT-aligned (0.15 pts)")
            total_score += 0.15
        elif not title_centered:
            print(f"FAIL: Component 3 — Title is not CENTER-aligned")
        else:
            print(f"FAIL: Component 3 — No content box is LEFT-aligned (compound check fails)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
