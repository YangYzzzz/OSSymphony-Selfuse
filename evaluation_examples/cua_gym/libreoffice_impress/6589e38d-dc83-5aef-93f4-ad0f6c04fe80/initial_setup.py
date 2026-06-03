"""
Initial Setup: Set text alignment on slide 4 of a psychology lecture presentation
Task ID: impress_teach_016
Domain: libreoffice_impress

Creates a 6-slide psychology lecture presentation. Slide 4 has a title and two
content text boxes, all center-aligned.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, alignment=PP_ALIGN.CENTER, font_color=None,
                font_name="Arial"):
    """Helper to add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_color:
        run.font.color.rgb = font_color
    return txBox


def add_multi_paragraph_textbox(slide, left, top, width, height, paragraphs,
                                 font_size=16, alignment=PP_ALIGN.CENTER,
                                 font_name="Arial", font_color=None):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.alignment = alignment
        if p.runs:
            run = p.runs[0]
            run.font.name = font_name
            run.font.size = Pt(font_size)
            if font_color:
                run.font.color.rgb = font_color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ============================================================
    # Slide 1: Title Slide
    # ============================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide1, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Introduction to Cognitive Psychology",
                font_size=36, bold=True, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))
    add_textbox(slide1, Inches(1.5), Inches(3.5), Inches(7), Inches(1),
                "Professor Elena Vasquez  |  PSY 301  |  Spring 2026",
                font_size=20, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x55, 0x55, 0x55))
    add_textbox(slide1, Inches(2), Inches(5), Inches(6), Inches(0.6),
                "Department of Psychology, Westfield University",
                font_size=14, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x88, 0x88, 0x88))

    # ============================================================
    # Slide 2: Course Overview
    # ============================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(8.4), Inches(1),
                "Course Overview",
                font_size=30, bold=True, alignment=PP_ALIGN.LEFT,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))
    add_multi_paragraph_textbox(
        slide2, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
        [
            "This course examines the mental processes underlying human behavior.",
            "Topics include perception, attention, memory, language, and decision-making.",
            "We will explore both classic experiments and contemporary research findings.",
            "Assessment includes two exams, a research paper, and weekly discussion posts.",
            "Office hours: Tuesdays and Thursdays, 2:00 PM - 4:00 PM in Room 412.",
        ],
        font_size=16, alignment=PP_ALIGN.LEFT
    )

    # ============================================================
    # Slide 3: Historical Foundations
    # ============================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(8.4), Inches(1),
                "Historical Foundations",
                font_size=30, bold=True, alignment=PP_ALIGN.LEFT,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))
    add_multi_paragraph_textbox(
        slide3, Inches(0.8), Inches(1.8), Inches(4), Inches(4.5),
        [
            "Wilhelm Wundt established the first psychology lab in Leipzig (1879).",
            "William James published Principles of Psychology in 1890.",
            "The behaviorist revolution shifted focus to observable behavior.",
            "The cognitive revolution of the 1950s-60s restored interest in mental processes.",
        ],
        font_size=15, alignment=PP_ALIGN.LEFT
    )
    add_multi_paragraph_textbox(
        slide3, Inches(5.2), Inches(1.8), Inches(4.2), Inches(4.5),
        [
            "Key figures: Noam Chomsky, George Miller, Ulric Neisser.",
            "Information processing model became the dominant framework.",
            "Computers provided a powerful metaphor for human cognition.",
        ],
        font_size=15, alignment=PP_ALIGN.LEFT
    )

    # ============================================================
    # Slide 4: Memory Systems (TARGET SLIDE)
    # All text boxes are CENTER-ALIGNED in initial state
    # ============================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    # Title - center-aligned
    add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(8.4), Inches(1),
                "Memory Systems and Encoding",
                font_size=30, bold=True, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))

    # Content text box 1 - center-aligned (task requires changing to left)
    add_multi_paragraph_textbox(
        slide4, Inches(0.8), Inches(1.7), Inches(4), Inches(4.5),
        [
            "Sensory memory retains brief impressions of stimuli.",
            "Short-term memory holds about 7 items for 20-30 seconds.",
            "Working memory actively manipulates information during tasks.",
            "Long-term memory stores information for extended periods.",
            "Encoding specificity principle links retrieval to encoding context.",
        ],
        font_size=15, alignment=PP_ALIGN.CENTER
    )

    # Content text box 2 - center-aligned (task requires changing to left)
    add_multi_paragraph_textbox(
        slide4, Inches(5.2), Inches(1.7), Inches(4.2), Inches(4.5),
        [
            "Elaborative rehearsal improves long-term retention.",
            "The spacing effect demonstrates distributed practice benefits.",
            "Mnemonics leverage existing knowledge to encode new information.",
            "Retrieval practice strengthens memory traces more than re-reading.",
        ],
        font_size=15, alignment=PP_ALIGN.CENTER
    )

    # ============================================================
    # Slide 5: Attention and Perception
    # ============================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(8.4), Inches(1),
                "Attention and Perception",
                font_size=30, bold=True, alignment=PP_ALIGN.LEFT,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))
    add_multi_paragraph_textbox(
        slide5, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
        [
            "Selective attention filters incoming sensory information.",
            "The cocktail party effect demonstrates automatic name detection.",
            "Change blindness reveals limitations in our visual awareness.",
            "Top-down processing uses expectations to interpret stimuli.",
            "Bottom-up processing builds perception from raw sensory data.",
            "Gestalt principles explain how we organize visual elements.",
        ],
        font_size=16, alignment=PP_ALIGN.LEFT
    )

    # ============================================================
    # Slide 6: Upcoming Topics
    # ============================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(8.4), Inches(1),
                "Next Week: Language and Thought",
                font_size=30, bold=True, alignment=PP_ALIGN.CENTER,
                font_color=RGBColor(0x1A, 0x3C, 0x6E))
    add_multi_paragraph_textbox(
        slide6, Inches(1.5), Inches(2.0), Inches(7), Inches(4),
        [
            "Reading: Chapters 8-9 in Goldstein (2024)",
            "Discussion post due Wednesday at 11:59 PM",
            "Research paper topic proposal due Friday",
            "Optional: Watch the TED Talk by Steven Pinker on language",
        ],
        font_size=18, alignment=PP_ALIGN.CENTER
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
