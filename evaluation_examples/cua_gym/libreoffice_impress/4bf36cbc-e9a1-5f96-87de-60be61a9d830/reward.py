"""
Reward Script: Animated countdown sequence on slide 1 (3, 2, 1) with entrance/exit animations
Task ID: impress_gf1_036
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Timing/animation element exists on slide 1
  Component 2 (0.30): All 3 text boxes ('3','2','1') have entrance animations
  Component 3 (0.25): All 3 text boxes have exit animations
  Component 4 (0.15): Correct countdown order (3 first, then 2, then 1)
  Component 5 (0.15): 1-second delay between entrance and exit for each number
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_036'

P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS = {'p': P_NS, 'a': A_NS}


def get_shape_text_map(pptx_path):
    """Map shape spid -> text content from slide 1."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    spid_text = {}
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            # Get the spid from the underlying XML element
            cNvPr_elems = shape._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}../../*')
            # Use a more direct approach: iterate through descendants
            for child in shape._element.iter():
                local = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                if local == 'cNvPr':
                    sid = child.get('id')
                    if sid:
                        spid_text[sid] = shape.text.strip()
                    break
    return spid_text


def parse_animations(pptx_path):
    """Parse animation XML from slide 1. Returns structured animation data."""
    result = {
        'has_timing': False,
        'entrance_anims': {},   # spid -> list of entrance anim nodes
        'exit_anims': {},       # spid -> list of exit anim nodes
        'anim_sequence': [],    # ordered list of (spid, class) tuples
        'delays': {},           # (spid, 'exit') -> delay in ms before exit
    }

    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            with zf.open('ppt/slides/slide1.xml') as f:
                root = ET.parse(f).getroot()
        except KeyError:
            return result

    timing = root.find(f'.//{{{P_NS}}}timing')
    if timing is None:
        return result
    result['has_timing'] = True

    # Find all animation nodes with presetClass (entrance/exit effects)
    # Walk through the main sequence to get order and properties
    main_seq = timing.find(f'.//{{{P_NS}}}seq')
    if main_seq is None:
        return result

    main_ctn = main_seq.find(f'{{{P_NS}}}cTn')
    if main_ctn is None:
        return result

    child_list = main_ctn.find(f'{{{P_NS}}}childTnLst')
    if child_list is None:
        return result

    # The main sequence has one top-level <par> containing all animation groups
    top_par = child_list.find(f'{{{P_NS}}}par')
    if top_par is None:
        return result

    top_ctn = top_par.find(f'{{{P_NS}}}cTn')
    if top_ctn is None:
        return result

    top_children = top_ctn.find(f'{{{P_NS}}}childTnLst')
    if top_children is None:
        return result

    # Each child <par> is an animation group (with a delay)
    for par in top_children.findall(f'{{{P_NS}}}par'):
        ctn = par.find(f'{{{P_NS}}}cTn')
        if ctn is None:
            continue

        # Get delay from stCondLst
        delay = 0
        st_cond = ctn.find(f'{{{P_NS}}}stCondLst/{{{P_NS}}}cond')
        if st_cond is not None:
            delay_str = st_cond.get('delay', '0')
            try:
                delay = int(delay_str)
            except ValueError:
                delay = 0

        # Find the animation effect inside
        inner_children = ctn.find(f'{{{P_NS}}}childTnLst')
        if inner_children is None:
            continue

        for inner_par in inner_children.findall(f'{{{P_NS}}}par'):
            inner_ctn = inner_par.find(f'{{{P_NS}}}cTn')
            if inner_ctn is None:
                continue

            preset_class = inner_ctn.get('presetClass', '')
            preset_id = inner_ctn.get('presetID', '')

            # Find target shape spid
            spid = None
            for spTgt in inner_ctn.iter(f'{{{P_NS}}}spTgt'):
                spid = spTgt.get('spid')
                if spid:
                    break

            if spid and preset_class:
                result['anim_sequence'].append((spid, preset_class, delay))

                if preset_class == 'entr':
                    result['entrance_anims'].setdefault(spid, []).append({
                        'presetID': preset_id,
                        'delay': delay,
                    })
                elif preset_class == 'exit':
                    result['exit_anims'].setdefault(spid, []).append({
                        'presetID': preset_id,
                        'delay': delay,
                    })
                    result['delays'][(spid, 'exit')] = delay

    return result


def verify_task(file_path):
    """
    Verify animated countdown sequence task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Get shape text map for identifying which shapes are '3', '2', '1'
    try:
        spid_text = get_shape_text_map(file_path)
        print(f"INFO: Shape text map: {spid_text}")
    except Exception as e:
        print(f"ERROR: Cannot parse shapes: {e}")
        spid_text = {}

    # Parse animations
    try:
        anims = parse_animations(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot parse animations: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Timing/animation element exists on slide 1 (0.15 points)
    try:
        if anims['has_timing']:
            print(f"PASS: Component 1 - Timing element exists on slide 1 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 - No timing/animation element found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: All 3 text boxes have entrance animations (0.30 points)
    # Award partial: 0.10 per text box with entrance animation
    try:
        entrance_count = 0
        target_spids = []
        for spid, text in spid_text.items():
            if text in ('3', '2', '1'):
                target_spids.append((spid, text))
                if spid in anims['entrance_anims']:
                    entrance_count += 1
                    print(f"  PASS: Text '{text}' (spid={spid}) has entrance animation")
                else:
                    print(f"  FAIL: Text '{text}' (spid={spid}) missing entrance animation")

        if entrance_count == 3:
            print(f"PASS: Component 2 - All 3 text boxes have entrance animations (0.30 pts)")
            total_score += 0.30
        elif entrance_count > 0:
            partial = entrance_count * 0.10
            print(f"PARTIAL: Component 2 - {entrance_count}/3 text boxes have entrance animations ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No text boxes have entrance animations")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: All 3 text boxes have exit animations (0.25 points)
    # Award partial: ~0.083 per text box with exit animation
    try:
        exit_count = 0
        for spid, text in target_spids:
            if spid in anims['exit_anims']:
                exit_count += 1
                print(f"  PASS: Text '{text}' (spid={spid}) has exit animation")
            else:
                print(f"  FAIL: Text '{text}' (spid={spid}) missing exit animation")

        if exit_count == 3:
            print(f"PASS: Component 3 - All 3 text boxes have exit animations (0.25 pts)")
            total_score += 0.25
        elif exit_count > 0:
            partial = round(exit_count * 0.25 / 3, 2)
            print(f"PARTIAL: Component 3 - {exit_count}/3 text boxes have exit animations ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No text boxes have exit animations")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Correct countdown order: '3' animates first, then '2', then '1' (0.15 points)
    try:
        # Build ordered list of entrance animations by their position in sequence
        entrance_order = []
        for spid, cls, delay in anims['anim_sequence']:
            if cls == 'entr' and spid in spid_text:
                entrance_order.append(spid_text[spid])

        print(f"  INFO: Entrance animation order: {entrance_order}")

        if entrance_order == ['3', '2', '1']:
            print(f"PASS: Component 4 - Correct countdown order 3->2->1 (0.15 pts)")
            total_score += 0.15
        elif set(entrance_order) == {'3', '2', '1'}:
            print(f"FAIL: Component 4 - All 3 numbers animated but wrong order: {entrance_order}")
        else:
            print(f"FAIL: Component 4 - Incomplete or missing countdown sequence: {entrance_order}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: 1-second (1000ms) delay between entrance and exit for each number (0.15 points)
    # Award partial: 0.05 per correct delay
    try:
        correct_delays = 0
        for spid, text in target_spids:
            key = (spid, 'exit')
            if key in anims['delays']:
                delay_ms = anims['delays'][key]
                # Allow some tolerance: 800-1200ms
                if 800 <= delay_ms <= 1200:
                    correct_delays += 1
                    print(f"  PASS: Text '{text}' exit delay = {delay_ms}ms (within tolerance)")
                else:
                    print(f"  FAIL: Text '{text}' exit delay = {delay_ms}ms (expected ~1000ms)")
            else:
                print(f"  FAIL: Text '{text}' has no exit delay info")

        if correct_delays == 3:
            print(f"PASS: Component 5 - All 3 numbers have ~1 second display time (0.15 pts)")
            total_score += 0.15
        elif correct_delays > 0:
            partial = round(correct_delays * 0.05, 2)
            print(f"PARTIAL: Component 5 - {correct_delays}/3 correct delays ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - No correct display delays found")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main execution
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
