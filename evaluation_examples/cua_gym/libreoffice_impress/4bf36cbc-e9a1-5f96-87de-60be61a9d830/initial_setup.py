"""
Initial Setup: Create Event_Opening.pptx with 6 slides, slide 1 has black background
and three centered text boxes with '3', '2', '1' in white 120pt bold.
Task ID: impress_gf1_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_centered_number_textbox(slide, number_text, prs):
    """Add a centered large number text box on the slide."""
    # Text box size
    box_w = Inches(4)
    box_h = Inches(2)
    # Center on slide
    left = (prs.slide_width - box_w) // 2
    top = (prs.slide_height - box_h) // 2
    txBox = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number_text
    run.font.size = Pt(120)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"
    return txBox


def create_initial():
    prs = Presentation()
    # Standard 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Countdown slide (black background, 3 number text boxes) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only or Blank
    # Black background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)

    # Remove any default placeholder shapes from the layout
    for ph in list(slide1.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # Add three centered number text boxes: 3, 2, 1
    # Slightly offset vertically so they overlap (as countdown numbers that appear/disappear)
    add_centered_number_textbox(slide1, "3", prs)
    add_centered_number_textbox(slide1, "2", prs)
    add_centered_number_textbox(slide1, "1", prs)

    # --- Slide 2: Event Title ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for ph in list(slide2.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "TechForward 2025"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Annual Innovation Summit"
    run2.font.size = Pt(28)
    run2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    run2.font.name = "Arial"

    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "March 22-24, 2025 | San Francisco Convention Center"
    run3.font.size = Pt(18)
    run3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run3.font.name = "Arial"

    # --- Slide 3: Agenda ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide3.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf3 = title3.text_frame
    p = tf3.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Event Agenda"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)
    run.font.name = "Arial"

    agenda_items = [
        ("9:00 AM", "Registration & Welcome Coffee"),
        ("10:00 AM", "Keynote: The Future of AI in Enterprise"),
        ("11:30 AM", "Panel: Sustainable Technology Practices"),
        ("12:30 PM", "Networking Lunch"),
        ("2:00 PM", "Workshop: Building Resilient Systems"),
        ("4:00 PM", "Closing Remarks & Awards Ceremony"),
    ]
    body3 = slide3.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf_body = body3.text_frame
    tf_body.word_wrap = True
    for i, (time_str, desc) in enumerate(agenda_items):
        if i == 0:
            p = tf_body.paragraphs[0]
        else:
            p = tf_body.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run_time = p.add_run()
        run_time.text = f"{time_str}  -  "
        run_time.font.size = Pt(20)
        run_time.font.bold = True
        run_time.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run_time.font.name = "Arial"
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(20)
        run_desc.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        run_desc.font.name = "Arial"

    # --- Slide 4: Speakers ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide4.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    title4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf4 = title4.text_frame
    p = tf4.paragraphs[0]
    run = p.add_run()
    run.text = "Featured Speakers"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)
    run.font.name = "Arial"

    speakers = [
        ("Dr. Maya Patel", "Chief AI Officer, NovaTech", "Leading expert in enterprise machine learning"),
        ("James Chen", "VP of Engineering, GreenScale", "Pioneer in sustainable computing infrastructure"),
        ("Sarah Williams", "CEO, FutureMinds", "Award-winning entrepreneur and keynote speaker"),
    ]
    for j, (name, role, bio) in enumerate(speakers):
        y_offset = Inches(1.8 + j * 1.8)
        sp_box = slide4.shapes.add_textbox(Inches(1.5), y_offset, Inches(10), Inches(1.5))
        tf_sp = sp_box.text_frame
        tf_sp.word_wrap = True
        p_name = tf_sp.paragraphs[0]
        run_n = p_name.add_run()
        run_n.text = name
        run_n.font.size = Pt(24)
        run_n.font.bold = True
        run_n.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        run_n.font.name = "Arial"
        p_role = tf_sp.add_paragraph()
        run_r = p_role.add_run()
        run_r.text = role
        run_r.font.size = Pt(18)
        run_r.font.italic = True
        run_r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        run_r.font.name = "Arial"
        p_bio = tf_sp.add_paragraph()
        run_b = p_bio.add_run()
        run_b.text = bio
        run_b.font.size = Pt(16)
        run_b.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        run_b.font.name = "Arial"

    # --- Slide 5: Venue Info ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in list(slide5.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    title5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
    tf5 = title5.text_frame
    p = tf5.paragraphs[0]
    run = p.add_run()
    run.text = "Venue & Logistics"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x6E)
    run.font.name = "Arial"

    info_box = slide5.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(4))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    venue_lines = [
        "San Francisco Convention Center - Hall B",
        "747 Howard Street, San Francisco, CA 94103",
        "",
        "Wi-Fi: TechForward2025  |  Password: innovate!25",
        "Parking: Lot C (validated for attendees)",
        "Emergency Contact: (415) 555-0198",
    ]
    for i, line in enumerate(venue_lines):
        if i == 0:
            p = tf_info.paragraphs[0]
        else:
            p = tf_info.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        run.font.name = "Arial"

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    for ph in list(slide6.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    ty_box = slide6.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(3))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    p = tf_ty.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You!"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    p2 = tf_ty.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "See you next year at TechForward 2026!"
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    run2.font.name = "Arial"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
