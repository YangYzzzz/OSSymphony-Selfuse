"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a 7×5 empty table at the caret position.
Generated: 2025-10-17 08:41:23
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
import os

def verify_table_task(file_path: str) -> float:
    """
    Verify that the presentation at `file_path` contains an empty 7×5 table
    (7 rows × 5 columns) inserted at the caret position.

    Scoring (progressive):
      • 0.0  – No relevant progress / file cannot be opened
      • 0.1  – Presentation contains *some* table, but not 7×5
      • 0.7  – Presentation contains a table whose grid matches 7×5
      • 1.0  – That table exists AND every single cell is empty

    Returns a float reward between 0.0 and 1.0 and prints detailed
    diagnostics plus a final line in form "REWARD: X.X".
    """

    print(f"Verifying PPTX file: {file_path}")
    score = 0.0
    max_score = 1.0

    # ----------- Basic file checks (no points awarded) ------------
    if not os.path.exists(file_path):
        print("✗ File does not exist.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ----------- Core verification logic --------------------------
    found_any_table = False          # Was *any* table detected?
    found_correct_table = False      # Is there a 7×5 (or 5×7) table?
    correct_table_empty = False      # Are all cells in that table empty?

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, 'has_table', False):
                continue  # Skip non-table shapes

            found_any_table = True
            table = shape.table
            rows, cols = len(table.rows), len(table.columns)
            print(f"Found table on slide {slide_idx}, shape {shape_idx} with dimensions {rows}×{cols}")

            # Accept either orientation: 7 rows × 5 columns OR 5 rows × 7 columns.
            if (rows == 7 and cols == 5) or (rows == 5 and cols == 7):
                found_correct_table = True

                # Check if every cell is empty (ignoring whitespace)
                all_cells_empty = True
                for r in table.rows:
                    for cell in r.cells:
                        if cell.text and cell.text.strip():
                            all_cells_empty = False
                            break
                    if not all_cells_empty:
                        break
                correct_table_empty = all_cells_empty

    # ------------------ Scoring -----------------------------------
    if found_any_table:
        print("✓ At least one table detected in presentation")

    if found_correct_table:
        print("✓ Found table with correct 7×5 configuration (0.7 points)")
        score += 0.7
        if correct_table_empty:
            print("✓ All table cells are empty (0.3 points)")
            score += 0.3
        else:
            print("✗ Table cells are not all empty (0.0 extra points)")
    else:
        if found_any_table:
            print("✗ No table with correct dimensions (0.1 partial credit)")
            score += 0.1
        else:
            print("✗ No tables found in presentation (0.0 points)")

    # Ensure score never exceeds 1.0
    final_score = min(score, max_score)

    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ----------------- Script execution when run directly -------------
if __name__ == "__main__":
    verify_table_task('/home/user/insert_a_75_empty_table_at_the_caret_position.pptx')
