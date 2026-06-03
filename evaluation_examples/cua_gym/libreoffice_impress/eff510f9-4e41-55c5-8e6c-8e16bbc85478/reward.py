"""
Reward Script: Insert a waterfall chart on slide 5 showing profit changes from revenue through cost deductions.
Task ID: impress_tct_069
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Chart exists on slide 5
  Component 2 (0.25): Waterfall categories are correct (Revenue, COGS, OpEx, Tax, Net Profit)
  Component 3 (0.25): Chart data values represent correct waterfall structure
  Component 4 (0.20): Color coding — green for positive, red for negative, invisible base series
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_069'

# Expected waterfall categories (normalized lowercase for flexible matching)
EXPECTED_CATEGORIES_NORMALIZED = ['revenue', 'cogs', 'operating expenses', 'tax', 'net profit']

# Expected data structure for 3-series waterfall (base, positive, negative)
# Revenue=19.5M, COGS=6.8M, OpEx=6.3M, Tax=1.6M, Net Profit=4.8M
EXPECTED_REVENUE = 19_500_000
EXPECTED_COGS = 6_800_000
EXPECTED_OPEX = 6_300_000
EXPECTED_TAX = 1_600_000
EXPECTED_NET_PROFIT = 4_800_000


def normalize_category(cat):
    """Normalize category text for comparison (lowercase, strip whitespace/newlines)."""
    return ' '.join(cat.lower().strip().split())


def find_chart_on_slide(slide):
    """Find a chart shape on the given slide. Returns shape or None."""
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Component 1: Chart exists on slide 5 (0.30 points)
    try:
        chart_shape = find_chart_on_slide(slide5)
        if chart_shape is not None:
            chart = chart_shape.chart
            print(f"PASS: Component 1 — Chart found on slide 5, type={chart.chart_type} (0.30 pts)")
            total_score += 0.30
        else:
            print("FAIL: Component 1 — No chart found on slide 5")
            # No chart means no further checks possible
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Waterfall categories are correct (0.25 points)
    try:
        plot = chart.plots[0]
        categories = [normalize_category(str(c)) for c in plot.categories]
        print(f"  Found categories: {categories}")

        # Check that all expected categories are present (order-sensitive)
        matches = 0
        for i, expected in enumerate(EXPECTED_CATEGORIES_NORMALIZED):
            if i < len(categories):
                # Allow partial match (e.g., "operating\nexpenses" -> "operating expenses")
                if expected in categories[i] or categories[i] in expected:
                    matches += 1
                else:
                    print(f"  Category mismatch at index {i}: expected '{expected}', got '{categories[i]}'")

        if matches == len(EXPECTED_CATEGORIES_NORMALIZED) and len(categories) == len(EXPECTED_CATEGORIES_NORMALIZED):
            print(f"PASS: Component 2 — All 5 waterfall categories correct (0.25 pts)")
            total_score += 0.25
        elif matches >= 3:
            partial = 0.25 * (matches / len(EXPECTED_CATEGORIES_NORMALIZED))
            print(f"PARTIAL: Component 2 — {matches}/5 categories match ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {matches}/5 categories match")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart data values represent correct waterfall structure (0.25 points)
    try:
        plot = chart.plots[0]
        num_series = len(plot.series)
        print(f"  Number of series: {num_series}")

        if num_series < 2:
            print("FAIL: Component 3 — Need at least 2 series for waterfall, found {num_series}")
        else:
            # For a waterfall chart, we expect:
            # - Revenue starts at 0 base, has a tall positive bar (19.5M)
            # - Intermediate steps show negative deductions (COGS, OpEx, Tax)
            # - Net Profit is the final positive bar (4.8M)
            #
            # Collect all series values
            all_values = []
            for si in range(num_series):
                vals = list(plot.series[si].values)
                all_values.append(vals)
                print(f"  Series {si} values: {vals}")

            # Check if the data forms a valid waterfall:
            # Sum of all series at each category should give the correct cumulative values
            num_cats = len(list(plot.categories))

            # Verify key financial values appear in the data
            # Flatten all values to check presence
            flat_values = set()
            for sv in all_values:
                for v in sv:
                    if v is not None and v != 0:
                        flat_values.add(v)

            value_checks = 0
            total_value_checks = 5

            # Check Revenue (19.5M)
            if EXPECTED_REVENUE in flat_values:
                value_checks += 1
                print(f"  Revenue {EXPECTED_REVENUE} found in chart data")
            else:
                print(f"  Revenue {EXPECTED_REVENUE} NOT found in chart data")

            # Check COGS (6.8M)
            if EXPECTED_COGS in flat_values:
                value_checks += 1
                print(f"  COGS {EXPECTED_COGS} found in chart data")
            else:
                print(f"  COGS {EXPECTED_COGS} NOT found in chart data")

            # Check OpEx (6.3M)
            if EXPECTED_OPEX in flat_values:
                value_checks += 1
                print(f"  OpEx {EXPECTED_OPEX} found in chart data")
            else:
                print(f"  OpEx {EXPECTED_OPEX} NOT found in chart data")

            # Check Tax (1.6M)
            if EXPECTED_TAX in flat_values:
                value_checks += 1
                print(f"  Tax {EXPECTED_TAX} found in chart data")
            else:
                print(f"  Tax {EXPECTED_TAX} NOT found in chart data")

            # Check Net Profit (4.8M)
            if EXPECTED_NET_PROFIT in flat_values:
                value_checks += 1
                print(f"  Net Profit {EXPECTED_NET_PROFIT} found in chart data")
            else:
                print(f"  Net Profit {EXPECTED_NET_PROFIT} NOT found in chart data")

            if value_checks >= 4:
                print(f"PASS: Component 3 — {value_checks}/5 key financial values present (0.25 pts)")
                total_score += 0.25
            elif value_checks >= 2:
                partial = 0.25 * (value_checks / total_value_checks)
                print(f"PARTIAL: Component 3 — {value_checks}/5 values present ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 — Only {value_checks}/5 key values found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Color coding — green for positive, red for negative, invisible base (0.20 points)
    try:
        plot = chart.plots[0]
        num_series = len(plot.series)
        color_score = 0.0

        # We need to identify which series plays which role
        # In a waterfall: one series is invisible base, one is positive (green), one is negative (red)
        has_invisible_base = False
        has_green_series = False
        has_red_series = False

        for si in range(num_series):
            series = plot.series[si]
            fill = series.format.fill
            fill_type = fill.type

            if fill_type is not None and fill_type == 5:  # BACKGROUND = invisible
                has_invisible_base = True
                print(f"  Series {si}: invisible base (BACKGROUND fill)")
            elif fill_type is not None and fill_type == 1:  # SOLID
                try:
                    rgb = fill.fore_color.rgb
                    rgb_str = str(rgb).upper()
                    r_val = int(rgb_str[0:2], 16)
                    g_val = int(rgb_str[2:4], 16)
                    b_val = int(rgb_str[4:6], 16)

                    # Green-ish: green channel dominant
                    if g_val > r_val and g_val > b_val:
                        has_green_series = True
                        print(f"  Series {si}: green ({rgb_str})")
                    # Red-ish: red channel dominant
                    elif r_val > g_val and r_val > b_val:
                        has_red_series = True
                        print(f"  Series {si}: red ({rgb_str})")
                    else:
                        print(f"  Series {si}: other color ({rgb_str})")
                except Exception as e:
                    print(f"  Series {si}: could not read color: {e}")
            else:
                print(f"  Series {si}: fill_type={fill_type}")

        # Award sub-points
        sub_checks = 0
        if has_invisible_base:
            sub_checks += 1
        if has_green_series:
            sub_checks += 1
        if has_red_series:
            sub_checks += 1

        if sub_checks == 3:
            print(f"PASS: Component 4 — Invisible base + green positive + red negative (0.20 pts)")
            total_score += 0.20
        elif sub_checks >= 1:
            partial = 0.20 * (sub_checks / 3)
            print(f"PARTIAL: Component 4 — {sub_checks}/3 color checks pass ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No waterfall color coding detected")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
