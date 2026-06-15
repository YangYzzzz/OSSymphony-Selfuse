"""
Initial Setup: Create a 6-slide presentation with slide 5 titled 'Profit Breakdown' but NO chart.
Task ID: impress_tct_069
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_069'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===================== Slide 1: Title Slide =====================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                "Meridian Technologies Inc.", font_size=36, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2),
                "Annual Financial Performance Report - FY 2024",
                font_size=22, color=RGBColor(0x4A, 0x4A, 0x4A),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.8),
                "Prepared by: Finance Department  |  Date: March 15, 2025",
                font_size=14, color=RGBColor(0x80, 0x80, 0x80),
                alignment=PP_ALIGN.CENTER)

    # ===================== Slide 2: Revenue Overview =====================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Revenue Overview", font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))

    # Revenue data table-like content
    revenue_lines = [
        "Product Sales:           $12,450,000    (+8.2% YoY)",
        "Service Revenue:          $4,320,000    (+12.5% YoY)",
        "Licensing & Royalties:    $1,890,000    (+3.1% YoY)",
        "Consulting Fees:            $840,000    (-2.4% YoY)",
        "",
        "Total Revenue:           $19,500,000    (+7.8% YoY)",
    ]
    y = Inches(1.6)
    for line in revenue_lines:
        if line:
            add_textbox(slide2, Inches(1.2), y, Inches(10), Inches(0.5),
                        line, font_size=16, color=RGBColor(0x33, 0x33, 0x33))
        y += Inches(0.55)

    add_textbox(slide2, Inches(1.2), Inches(5.8), Inches(10), Inches(0.8),
                "Key driver: Strong Q3/Q4 enterprise software sales pushed product revenue above forecast.",
                font_size=13, color=RGBColor(0x66, 0x66, 0x66))

    # ===================== Slide 3: Cost Structure =====================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Cost Structure", font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))

    cost_lines = [
        "Cost of Goods Sold (COGS):",
        "  Direct Materials:       $3,120,000",
        "  Direct Labor:           $2,480,000",
        "  Manufacturing OH:       $1,200,000",
        "  Total COGS:             $6,800,000   (34.9% of revenue)",
        "",
        "Gross Profit:            $12,700,000   (65.1% margin)",
    ]
    y = Inches(1.6)
    for line in cost_lines:
        if line:
            is_total = "Total COGS" in line or "Gross Profit" in line
            add_textbox(slide3, Inches(1.2), y, Inches(10), Inches(0.5),
                        line, font_size=16, bold=is_total,
                        color=RGBColor(0x33, 0x33, 0x33))
        y += Inches(0.55)

    # ===================== Slide 4: Operating Expenses =====================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Operating Expenses", font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))

    opex_lines = [
        "Research & Development:   $2,800,000   (14.4% of revenue)",
        "Sales & Marketing:        $2,100,000   (10.8% of revenue)",
        "General & Administrative: $1,400,000    (7.2% of revenue)",
        "",
        "Total Operating Expenses: $6,300,000   (32.3% of revenue)",
        "",
        "Operating Income:         $6,400,000   (32.8% margin)",
    ]
    y = Inches(1.6)
    for line in opex_lines:
        if line:
            is_total = "Total Operating" in line or "Operating Income" in line
            add_textbox(slide4, Inches(1.2), y, Inches(10), Inches(0.5),
                        line, font_size=16, bold=is_total,
                        color=RGBColor(0x33, 0x33, 0x33))
        y += Inches(0.55)

    # ===================== Slide 5: Profit Breakdown (NO CHART) =====================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Profit Breakdown", font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))

    add_textbox(slide5, Inches(1.2), Inches(2.5), Inches(10), Inches(1),
                "A waterfall chart will be inserted here to visualize the flow from Revenue to Net Profit.",
                font_size=16, color=RGBColor(0x99, 0x99, 0x99),
                alignment=PP_ALIGN.CENTER)

    # ===================== Slide 6: Summary & Outlook =====================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
                "Summary & Outlook", font_size=30, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C))

    summary_lines = [
        "FY 2024 Highlights:",
        "  - Revenue grew 7.8% YoY to $19.5M, exceeding guidance of $18.8M",
        "  - Gross margin improved 1.2pp to 65.1% through supply chain optimization",
        "  - Operating margin held steady at 32.8% despite increased R&D spending",
        "  - Net profit reached $4.8M after $1.6M in tax provisions",
        "",
        "FY 2025 Outlook:",
        "  - Revenue target: $21.5M - $22.0M (+10-13% growth)",
        "  - Planned R&D investment increase to 16% of revenue",
        "  - Expected operating margin: 30-32%",
    ]
    y = Inches(1.6)
    for line in summary_lines:
        if line:
            is_header = line.endswith(":")
            add_textbox(slide6, Inches(1.2), y, Inches(10), Inches(0.5),
                        line, font_size=15 if not is_header else 17,
                        bold=is_header,
                        color=RGBColor(0x33, 0x33, 0x33))
        y += Inches(0.5)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
