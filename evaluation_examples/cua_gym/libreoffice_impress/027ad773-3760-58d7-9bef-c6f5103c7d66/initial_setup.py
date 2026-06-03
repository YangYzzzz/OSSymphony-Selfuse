"""
Initial Setup: Create a 7-slide tech proposal presentation with inconsistent fonts.
Task ID: impstruct_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# --- Inconsistent fonts/sizes to use across slides ---
TITLE_FONTS = [
    ("Arial", 36),
    ("Times New Roman", 44),
    ("Courier New", 32),
    ("DejaVu Sans", 38),
    ("Liberation Sans", 42),
    ("FreeSans", 34),
    ("Nimbus Roman", 40),
]

BODY_FONTS = [
    ("Arial", 16),
    ("Times New Roman", 12),
    ("Courier New", 18),
    ("DejaVu Sans", 11),
    ("Liberation Sans", 15),
    ("FreeSans", 13),
    ("Nimbus Roman", 16),
]

# --- Slide content for a realistic tech proposal ---
SLIDES = [
    {
        "title": "Cloud Migration Strategy Proposal",
        "body": (
            "Prepared by: Horizon Digital Solutions\n"
            "Date: March 2025\n"
            "Version 3.2 — Confidential\n"
            "Client: Meridian Financial Group"
        ),
    },
    {
        "title": "Executive Summary",
        "body": (
            "Meridian Financial Group currently operates 42 on-premise servers\n"
            "across three data centers in Chicago, Denver, and Atlanta.\n"
            "Annual infrastructure costs exceed $2.8M with 99.1% uptime.\n"
            "This proposal outlines a phased migration to AWS/Azure hybrid\n"
            "cloud targeting 99.95% uptime and 35% cost reduction."
        ),
    },
    {
        "title": "Current Infrastructure Assessment",
        "body": (
            "Primary data center (Chicago): 18 servers, 4.2 PB storage\n"
            "Secondary (Denver): 14 servers, 2.8 PB storage\n"
            "Disaster recovery (Atlanta): 10 servers, 1.5 PB storage\n"
            "Peak utilization: 73% CPU, 81% memory during Q4 processing\n"
            "Hardware refresh cycle: 60% due for replacement by Q2 2026"
        ),
    },
    {
        "title": "Proposed Architecture",
        "body": (
            "Phase 1: Migrate dev/test workloads to AWS (Months 1-3)\n"
            "Phase 2: Production web tier to Azure App Services (Months 4-6)\n"
            "Phase 3: Database migration with zero-downtime cutover (Months 7-9)\n"
            "Phase 4: Legacy system containerization (Months 10-12)\n"
            "Retained on-prem: HSM modules and regulatory-bound data stores"
        ),
    },
    {
        "title": "Cost-Benefit Analysis",
        "body": (
            "Current annual spend: $2,840,000\n"
            "Projected cloud annual spend: $1,846,000 (35% savings)\n"
            "One-time migration cost: $620,000\n"
            "Break-even point: Month 18 post-migration\n"
            "5-year total savings projection: $4,350,000"
        ),
    },
    {
        "title": "Risk Mitigation Framework",
        "body": (
            "Data sovereignty: All regulated data remains in US-East regions\n"
            "Vendor lock-in: Kubernetes abstraction layer across providers\n"
            "Rollback capability: Parallel systems for 90 days post-migration\n"
            "Security: Zero-trust architecture with CrowdStrike integration\n"
            "Compliance: SOC 2 Type II and PCI-DSS maintained throughout"
        ),
    },
    {
        "title": "Timeline and Next Steps",
        "body": (
            "Kick-off meeting: April 14, 2025\n"
            "Infrastructure audit completion: May 9, 2025\n"
            "Phase 1 go-live: July 21, 2025\n"
            "Full migration target: March 31, 2026\n"
            "Quarterly review cadence with Meridian leadership team"
        ),
    },
]


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_content in enumerate(SLIDES):
        # Use layout 0 (Title Slide) for first slide, layout 1 (Title+Content) for rest
        layout_idx = 0 if i == 0 else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # --- Set title with inconsistent font ---
        title_shape = slide.shapes.title
        title_shape.text = ""
        p = title_shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = slide_content["title"]
        tfont, tsize = TITLE_FONTS[i]
        run.font.name = tfont
        run.font.size = Pt(tsize)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # --- Set body with inconsistent font ---
        body_shape = slide.placeholders[1]
        body_shape.text = ""
        tf = body_shape.text_frame
        lines = slide_content["body"].split("\n")
        for j, line in enumerate(lines):
            if j == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            run = para.add_run()
            run.text = line
            bfont, bsize = BODY_FONTS[i]
            run.font.name = bfont
            run.font.size = Pt(bsize)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
