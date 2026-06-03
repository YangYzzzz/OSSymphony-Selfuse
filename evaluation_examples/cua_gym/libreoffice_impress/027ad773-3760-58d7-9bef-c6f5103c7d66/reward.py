"""
Reward Script: Change all title fonts to Liberation Serif 40pt and body text to Liberation Mono 14pt
Task ID: impstruct_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All title placeholders use Liberation Serif 40pt
  Component 2 (0.5): All body/subtitle placeholders use Liberation Mono 14pt
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impstruct_018'

# Title placeholder types (CENTER_TITLE=3, TITLE=1)
TITLE_PH_TYPES = {1, 3}
# Body placeholder types (SUBTITLE=4, OBJECT=7, BODY=6)
BODY_PH_TYPES = {4, 6, 7}

EXPECTED_TITLE_FONT = 'Liberation Serif'
EXPECTED_TITLE_SIZE = Pt(40)  # 508000 EMU
EXPECTED_BODY_FONT = 'Liberation Mono'
EXPECTED_BODY_SIZE = Pt(14)   # 177800 EMU


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: 7 slides
    if len(prs.slides) != 7:
        print(f"WARN: Expected 7 slides, found {len(prs.slides)}")

    # Component 1: All title placeholders use Liberation Serif 40pt (0.5 points)
    try:
        title_total = 0
        title_pass = 0
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'placeholder_format') or shape.placeholder_format is None:
                    continue
                ph_type = shape.placeholder_format.type
                if ph_type not in TITLE_PH_TYPES:
                    continue
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    for run in runs:
                        title_total += 1
                        font_ok = (run.font.name == EXPECTED_TITLE_FONT)
                        size_ok = (run.font.size == EXPECTED_TITLE_SIZE)
                        if font_ok and size_ok:
                            title_pass += 1
                        else:
                            print(f"FAIL: Slide {slide_idx+1} title run: "
                                  f"font={run.font.name} (expect {EXPECTED_TITLE_FONT}), "
                                  f"size={run.font.size} (expect {EXPECTED_TITLE_SIZE})")

        if title_total == 0:
            print("FAIL: Component 1 -- no title runs found")
        elif title_pass == title_total:
            print(f"PASS: Component 1 -- all {title_total} title runs are Liberation Serif 40pt (0.5 pts)")
            total_score += 0.5
        else:
            partial = 0.5 * (title_pass / title_total)
            print(f"PARTIAL: Component 1 -- {title_pass}/{title_total} title runs correct ({partial:.3f} pts)")
            if partial > 0:
                total_score += partial
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All body/subtitle placeholders use Liberation Mono 14pt (0.5 points)
    try:
        body_total = 0
        body_pass = 0
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not hasattr(shape, 'placeholder_format') or shape.placeholder_format is None:
                    continue
                ph_type = shape.placeholder_format.type
                if ph_type not in BODY_PH_TYPES:
                    continue
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    for run in runs:
                        body_total += 1
                        font_ok = (run.font.name == EXPECTED_BODY_FONT)
                        size_ok = (run.font.size == EXPECTED_BODY_SIZE)
                        if font_ok and size_ok:
                            body_pass += 1
                        else:
                            print(f"FAIL: Slide {slide_idx+1} body run: "
                                  f"font={run.font.name} (expect {EXPECTED_BODY_FONT}), "
                                  f"size={run.font.size} (expect {EXPECTED_BODY_SIZE})")

        if body_total == 0:
            print("FAIL: Component 2 -- no body runs found")
        elif body_pass == body_total:
            print(f"PASS: Component 2 -- all {body_total} body runs are Liberation Mono 14pt (0.5 pts)")
            total_score += 0.5
        else:
            partial = 0.5 * (body_pass / body_total)
            print(f"PARTIAL: Component 2 -- {body_pass}/{body_total} body runs correct ({partial:.3f} pts)")
            if partial > 0:
                total_score += partial
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
