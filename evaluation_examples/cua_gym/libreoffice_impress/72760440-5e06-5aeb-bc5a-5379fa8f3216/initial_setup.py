"""
Initial Setup: Create a 10-slide Research Lecture presentation (no bibliography)
Task ID: impress_teach_088
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_088'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(18)
        p.level = 0
    return slide


def add_blank_with_textbox(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Machine Learning in Climate Science",
                    "Dr. Elena Vasquez\nDepartment of Environmental Data Science\nStanford University")

    # Slide 2: Overview
    add_content_slide(prs, "Lecture Overview", [
        "Motivation: Why ML for climate modeling?",
        "Key ML architectures used in climate research",
        "Case studies: Temperature prediction and extreme events",
        "Current limitations and open challenges",
        "Future directions in hybrid modeling"
    ])

    # Slide 3: Background
    add_content_slide(prs, "Background: Climate Modeling Challenges", [
        "Global Climate Models (GCMs) require enormous compute resources",
        "Sub-grid processes remain poorly resolved at typical 50-100 km resolution",
        "Parameterization schemes introduce systematic biases",
        "Observational data is sparse over oceans and polar regions",
        "Ensemble runs are limited by computational cost"
    ])

    # Slide 4: ML Architectures
    add_content_slide(prs, "ML Architectures in Climate Science", [
        "Convolutional Neural Networks for spatial pattern recognition",
        "Long Short-Term Memory networks for temporal sequences",
        "Graph Neural Networks for irregular climate grids",
        "Transformer models for global teleconnection patterns",
        "Physics-Informed Neural Networks (PINNs) for conservation laws"
    ])

    # Slide 5: Case Study 1
    add_content_slide(prs, "Case Study: Regional Temperature Downscaling", [
        "Objective: 50 km GCM output to 5 km resolution",
        "Dataset: ERA5 reanalysis (1979-2023) + station observations",
        "Method: U-Net with topographic features as auxiliary input",
        "Results: 38% RMSE reduction vs. bilinear interpolation",
        "Validation against independent weather station records"
    ])

    # Slide 6: Case Study 2
    add_content_slide(prs, "Case Study: Extreme Event Detection", [
        "Tropical cyclone track prediction using ConvLSTM",
        "Heatwave probability estimation with Random Forests",
        "Drought onset detection from satellite imagery (MODIS + Sentinel)",
        "Flood mapping using Transformer-based segmentation",
        "Multi-hazard early warning system integration"
    ])

    # Slide 7: Limitations
    add_content_slide(prs, "Current Limitations", [
        "Black-box nature: difficulty interpreting learned representations",
        "Distribution shift: models trained on historical data may not generalize",
        "Data quality: inconsistent measurement standards across networks",
        "Physical consistency: ML outputs may violate conservation laws",
        "Computational cost of training large foundation models"
    ])

    # Slide 8: Future Directions
    add_content_slide(prs, "Future Directions", [
        "Foundation models pre-trained on multi-petabyte climate datasets",
        "Hybrid models combining physical equations with learned corrections",
        "Federated learning across international climate research centers",
        "Uncertainty quantification through Bayesian deep learning",
        "Real-time inference for operational weather forecasting"
    ])

    # Slide 9: Key Takeaways
    add_content_slide(prs, "Key Takeaways", [
        "ML complements, not replaces, physics-based climate models",
        "Downscaling and extreme event detection are high-impact applications",
        "Physical constraints improve generalization and trustworthiness",
        "Collaboration between climate scientists and ML researchers is essential",
        "Open datasets and reproducibility standards are critical for progress"
    ])

    # Slide 10: Questions
    add_blank_with_textbox(prs, "Questions & Discussion",
        "Thank you for your attention.\n\nContact: e.vasquez@stanford.edu\nLab: vasquez-lab.stanford.edu\n\nOffice Hours: Tuesdays 2-4 PM, Room 312 Green Earth Sciences")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
