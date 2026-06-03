"""
Reward Script: Add bibliography slide with APA citations and hanging indent
Task ID: impress_teach_088
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 11 exists (0.2 pts)
  Component 2: Title 'References' on slide 11 (0.2 pts)
  Component 3: 5 citation paragraphs present (0.2 pts)
  Component 4: Font size 14pt on all citations (0.2 pts)
  Component 5: Hanging indent 0.5 inch on all citations (0.2 pts)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_088'


def persist_app_state(domain):
    """Save any unsaved changes in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt, Emu
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: Slide 11 exists (0.2 points)
    # Initial has 10 slides, golden has 11. This checks the task-introduced change.
    try:
        if num_slides >= 11:
            print(f"PASS: Component 1 — Slide 11 exists ({num_slides} slides) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 1 — Expected >= 11 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: if no slide 11, remaining checks are impossible
    if num_slides < 11:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    slide11 = prs.slides[10]  # 0-indexed

    # Component 2: Title 'References' on slide 11 (0.2 points)
    # Initial has no slide 11 at all, so this only passes on golden.
    try:
        references_found = False
        for shape in slide11.shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text.lower() == 'references':
                    references_found = True
                    break
        if references_found:
            print(f"PASS: Component 2 — 'References' title found on slide 11 (0.2 pts)")
            total_score += 0.2
        else:
            all_texts = []
            for shape in slide11.shapes:
                if shape.has_text_frame:
                    all_texts.append(shape.text_frame.text.strip()[:60])
            print(f"FAIL: Component 2 — No 'References' title found. Shape texts: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Find the citations text box (the shape with multiple paragraphs, not the title)
    citations_shape = None
    try:
        for shape in slide11.shapes:
            if shape.has_text_frame:
                # The citations shape has multiple non-empty paragraphs and is not the title
                nonempty = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                if len(nonempty) >= 3 and shape.text_frame.text.strip().lower() != 'references':
                    citations_shape = shape
                    break
    except Exception as e:
        print(f"ERROR: Finding citations shape — {e}")

    if citations_shape is None:
        print("FAIL: No citations text box found on slide 11")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    citation_paras = [p for p in citations_shape.text_frame.paragraphs if p.text.strip()]
    print(f"INFO: Found citations shape with {len(citation_paras)} non-empty paragraphs")

    # Component 3: 5 citation paragraphs present (0.2 points)
    # Initial has no slide 11, so this only passes on golden.
    try:
        if len(citation_paras) == 5:
            print(f"PASS: Component 3 — Exactly 5 citation paragraphs found (0.2 pts)")
            total_score += 0.2
        elif len(citation_paras) >= 4:
            # Partial credit: 4 citations = 0.1
            partial = 0.1
            print(f"PARTIAL: Component 3 — Found {len(citation_paras)} citations (expected 5) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Expected 5 citations, found {len(citation_paras)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Font size 14pt (177800 EMU) on all citation runs (0.2 points)
    # Initial has no slide 11 with these properties.
    try:
        target_size = Pt(14)  # 177800 EMU
        total_runs = 0
        correct_runs = 0
        for para in citation_paras:
            for run in para.runs:
                if run.text.strip():
                    total_runs += 1
                    if run.font.size is not None and run.font.size == target_size:
                        correct_runs += 1
                    else:
                        print(f"  INFO: Run font size = {run.font.size} (expected {target_size})")

        if total_runs > 0 and correct_runs == total_runs:
            print(f"PASS: Component 4 — All {total_runs} runs are 14pt (0.2 pts)")
            total_score += 0.2
        elif total_runs > 0:
            ratio = correct_runs / total_runs
            partial = round(0.2 * ratio, 2)
            print(f"PARTIAL: Component 4 — {correct_runs}/{total_runs} runs are 14pt ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No text runs found in citations")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Hanging indent formatting (0.2 points)
    # Hanging indent: marL = 457200 (0.5 inch), indent = -457200 (negative = hanging)
    # 0.5 inch = 457200 EMU. Allow some tolerance (within 10%).
    try:
        expected_marL = 457200  # 0.5 inch in EMU
        expected_indent = -457200
        tolerance = 0.15  # 15% tolerance

        paras_checked = 0
        paras_correct = 0
        for para in citation_paras:
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                marL_str = pPr.get('marL')
                indent_str = pPr.get('indent')
                marL_val = int(marL_str) if marL_str else None
                indent_val = int(indent_str) if indent_str else None

                paras_checked += 1
                # Check marL is approximately 0.5 inch (positive)
                marL_ok = (marL_val is not None and
                           abs(marL_val - expected_marL) / expected_marL <= tolerance)
                # Check indent is approximately -0.5 inch (negative = hanging)
                indent_ok = (indent_val is not None and indent_val < 0 and
                             abs(indent_val - expected_indent) / abs(expected_indent) <= tolerance)

                if marL_ok and indent_ok:
                    paras_correct += 1
                else:
                    print(f"  INFO: Para indent: marL={marL_val}, indent={indent_val} "
                          f"(expected marL~{expected_marL}, indent~{expected_indent})")
            else:
                paras_checked += 1
                print(f"  INFO: Paragraph has no pPr element (no indent)")

        if paras_checked > 0 and paras_correct == paras_checked:
            print(f"PASS: Component 5 — All {paras_checked} citation paragraphs have hanging indent (0.2 pts)")
            total_score += 0.2
        elif paras_checked > 0 and paras_correct > 0:
            ratio = paras_correct / paras_checked
            partial = round(0.2 * ratio, 2)
            print(f"PARTIAL: Component 5 — {paras_correct}/{paras_checked} paragraphs have correct indent ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No paragraphs have hanging indent formatting")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
