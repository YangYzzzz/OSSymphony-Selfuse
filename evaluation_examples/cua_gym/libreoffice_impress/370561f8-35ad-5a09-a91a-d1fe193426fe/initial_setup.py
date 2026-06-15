"""
Initial Setup: Create presentation with overlapping blue rectangle and red circle on slide 3
Task ID: impress_ndo_060
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_060'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Geometric Design Workshop"
    slide1.placeholders[1].text = "Visual Composition & Shape Operations\nQ2 2025 Training Series"

    # ── Slide 2: Agenda ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Session 1: Basic Shape Properties"
    items = [
        "Session 2: Shape Alignment & Distribution",
        "Session 3: Boolean Shape Operations (Union, Intersect, Subtract)",
        "Session 4: Advanced Composition Techniques",
        "Session 5: Practical Design Exercises",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # ── Slide 3: Two overlapping shapes (the task slide) ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box at the top
    title_box = slide3.shapes.add_textbox(Cm(2), Cm(0.5), Cm(20), Cm(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Shape Merge Exercise"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Blue rectangle: 8cm x 5cm, positioned center-left
    # Center of slide: ~12.7cm wide, ~9.525cm tall
    rect_left = Cm(7)
    rect_top = Cm(4)
    rect_width = Cm(8)
    rect_height = Cm(5)
    rect = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, rect_left, rect_top, rect_width, rect_height
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)  # #3498DB blue
    rect.line.color.rgb = RGBColor(0x2C, 0x80, 0xB8)
    rect.line.width = Pt(1)
    rect.name = "BlueRectangle"

    # Red circle: 5cm diameter, overlapping the rectangle
    # Position so it overlaps the right portion of the rectangle
    circle_diameter = Cm(5)
    circle_left = Cm(11)  # overlaps with rectangle (rect goes from 7 to 15cm)
    circle_top = Cm(4)    # same top as rectangle
    circle = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, circle_left, circle_top, circle_diameter, circle_diameter
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xE7, 0x4C, 0x3C)  # #E74C3C red
    circle.line.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    circle.line.width = Pt(1)
    circle.name = "RedCircle"

    # Add instruction text below the shapes
    instr_box = slide3.shapes.add_textbox(Cm(3), Cm(10), Cm(18), Cm(2))
    tf2 = instr_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Exercise: Merge the blue rectangle and red circle using the Union operation"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(14)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    # ── Slide 4: Additional content ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Boolean Shape Operations"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Union: Combines multiple shapes into one"
    ops = [
        "Intersect: Keeps only overlapping area",
        "Subtract: Removes one shape from another",
        "Fragment: Splits shapes at intersections",
        "Combine: Toggles overlap regions",
    ]
    for op in ops:
        p = body4.add_paragraph()
        p.text = op
        p.level = 0

    # ── Slide 5: Summary ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Takeaways"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Boolean operations enable complex shape design"
    takeaways = [
        "Union is ideal for creating composite icons and logos",
        "Always select shapes in the correct order for Subtract",
        "Use Fragment for creating puzzle-piece layouts",
    ]
    for t in takeaways:
        p = body5.add_paragraph()
        p.text = t
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
