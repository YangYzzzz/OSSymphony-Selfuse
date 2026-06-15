"""
Reward Script: Union merge of two shapes on slide 3 and recolor to #9B59B6
Task ID: impress_ndo_060
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Original BlueRectangle and RedCircle no longer exist as separate shapes
  Component 2 (0.30): A merged shape (freeform or custom geometry) exists on slide 3
  Component 3 (0.35): The merged shape has solid fill color #9B59B6
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_060'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"PRECONDITION FAIL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, so slide 3 is index 2
    shape_names = [s.name for s in slide3.shapes]
    shape_types = {s.name: s.shape_type for s in slide3.shapes}

    # Component 1: Original BlueRectangle and RedCircle no longer exist as separate shapes (0.35 points)
    # In initial_env, both exist. After union, they should be merged into one.
    try:
        blue_exists = 'BlueRectangle' in shape_names
        red_exists = 'RedCircle' in shape_names

        # Also check by type: count AUTO_SHAPE instances that could be the originals
        # The originals were AUTO_SHAPE type with specific fills
        auto_shapes_on_slide = [s for s in slide3.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        original_color_shapes = []
        for s in auto_shapes_on_slide:
            try:
                if s.fill.type is not None:
                    color = str(s.fill.fore_color.rgb).upper()
                    if color in ('3498DB', 'E74C3C'):
                        original_color_shapes.append(color)
            except Exception:
                pass

        if not blue_exists and not red_exists and len(original_color_shapes) == 0:
            print(f"PASS: Component 1 -- Original shapes removed (BlueRectangle and RedCircle gone) (0.35 pts)")
            total_score += 0.35
        else:
            details = []
            if blue_exists:
                details.append("BlueRectangle still present")
            if red_exists:
                details.append("RedCircle still present")
            if len(original_color_shapes) > 0:
                details.append("AUTO_SHAPE with original color still present")
            print(f"FAIL: Component 1 -- {'; '.join(details)}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: A merged shape exists on slide 3 (0.30 points)
    # After union, there should be a freeform/custom-geometry shape.
    # We check for FREEFORM type or a shape named 'MergedShape' or similar non-original shape.
    try:
        freeform_shapes = [s for s in slide3.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
        # Also check for custom geometry shapes that might result from union
        # MSO_SHAPE_TYPE values: FREEFORM=5, GROUP=6, AUTO_SHAPE=1
        # Union operations typically produce freeform shapes
        merged_candidates = freeform_shapes

        # Additionally, check for any non-original shape that could be the merged result
        # (in case the shape type differs from FREEFORM after union in LO)
        original_names = {'Title 1', 'TextBox 2', 'BlueRectangle', 'RedCircle', 'TextBox 5'}
        non_original = [s for s in slide3.shapes
                        if s.name not in original_names
                        and s.shape_type not in (MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.TEXT_BOX)]

        if len(merged_candidates) > 0:
            print(f"PASS: Component 2 -- Merged freeform shape found ({len(merged_candidates)} freeform(s)) (0.30 pts)")
            total_score += 0.30
        elif len(non_original) > 0:
            # Might be a different shape type from the union
            print(f"PASS: Component 2 -- Non-original shape found: {[s.name for s in non_original]} (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 -- No merged/freeform shape found on slide 3. Shapes: {shape_names}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: The merged shape has fill color #9B59B6 (0.35 points)
    # After union, the resulting shape should be filled with purple #9B59B6
    try:
        target_color = '9B59B6'
        matched_shape_name = None

        # Look for the merged shape (freeform or non-original) and check its fill
        candidates = [s for s in slide3.shapes if s.shape_type == MSO_SHAPE_TYPE.FREEFORM]
        if not candidates:
            # Fallback: any non-placeholder, non-textbox shape that isn't original
            original_names_set = {'Title 1', 'TextBox 2', 'BlueRectangle', 'RedCircle', 'TextBox 5'}
            candidates = [s for s in slide3.shapes
                          if s.name not in original_names_set
                          and s.shape_type not in (MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.TEXT_BOX)]

        for shape in candidates:
            try:
                if shape.fill.type is not None:
                    actual_color = str(shape.fill.fore_color.rgb).upper()
                    if actual_color == target_color:
                        matched_shape_name = shape.name
                        print(f"PASS: Component 3 -- Merged shape '{shape.name}' has fill #{actual_color} (0.35 pts)")
                        total_score += 0.35
                        break
                    else:
                        print(f"INFO: Shape '{shape.name}' has fill #{actual_color}, expected #{target_color}")
            except Exception as ex:
                print(f"INFO: Could not read fill of shape '{shape.name}': {ex}")

        if matched_shape_name is None and not candidates:
            print(f"FAIL: Component 3 -- No merged shape candidate found to check color")
        elif matched_shape_name is None:
            print(f"FAIL: Component 3 -- No shape with fill #{target_color} found among candidates")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 1)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
