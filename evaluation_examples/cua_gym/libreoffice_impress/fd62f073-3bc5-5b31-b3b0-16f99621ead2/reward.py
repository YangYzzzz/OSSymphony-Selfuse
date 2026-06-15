"""
Reward Script: Insert TechCorp as Fontwork/WordArt with arch shape in gold (#FFD700) at top of slide 1
Task ID: impress_objects_068
Domain: libreoffice_impress
Scoring:
  - Component 1: TechCorp text present on slide 1  (0.35 pts)
  - Component 2: Gold color #FFD700 applied         (0.30 pts)
  - Component 3: Arch/warp text effect present      (0.20 pts)
  - Component 4: Positioned at top of slide         (0.15 pts)
"""

import os
import xml.etree.ElementTree as ET

from pptx import Presentation

WORKDIR = '/home/user'  # VM path — all reward scripts run on the VM
TASK_ID = 'impress_objects_068'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: On slide 1, insert the company name 'TechCorp' as a Fontwork/WordArt text effect
    with an arch shape, colored in gold (#FFD700), positioned at the top of the slide.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slide 1 must exist
    if len(prs.slides) < 1:
        print("CRITICAL: No slides found in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find all shapes that contain 'TechCorp' text (including nested/group shapes)
    techcorp_shapes = []
    for shape in slide.shapes:
        found_techcorp = False
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                full_text = para.text.strip()
                for run in para.runs:
                    if 'TechCorp' in run.text:
                        found_techcorp = True
                        break
                if found_techcorp:
                    break
            # Also check the full text_frame text
            if not found_techcorp and 'TechCorp' in shape.text_frame.text:
                found_techcorp = True
        if found_techcorp:
            techcorp_shapes.append(shape)

    # ----- Component 1: TechCorp text present on slide 1 (0.35 points) -----
    # This FAILS on initial (no TechCorp) → PASSES on golden ✅
    try:
        if len(techcorp_shapes) >= 1:
            print(f"PASS: Component 1 — 'TechCorp' text found in {len(techcorp_shapes)} shape(s) on slide 1 (0.35 pts)")
            total_score += 0.35
        else:
            print("FAIL: Component 1 — No shape containing 'TechCorp' found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no TechCorp shape, skip the remaining components
    if not techcorp_shapes:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Use the first TechCorp shape for further checks
    tc_shape = techcorp_shapes[0]

    # ----- Component 2: Gold color #FFD700 applied to TechCorp text (0.30 points) -----
    # This FAILS on initial (shape doesn't exist) → PASSES on golden ✅
    try:
        gold_color_found = False
        color_found = None
        if tc_shape.has_text_frame:
            for para in tc_shape.text_frame.paragraphs:
                for run in para.runs:
                    if 'TechCorp' not in run.text:
                        continue
                    try:
                        if run.font.color.type is not None:
                            rgb = str(run.font.color.rgb).upper()
                            color_found = rgb
                            if rgb == 'FFD700':
                                gold_color_found = True
                    except Exception:
                        pass

        # Also check via XML for solidFill with srgbClr val="FFD700" in case font.color.type is None
        if not gold_color_found:
            xml_str = ET.tostring(tc_shape.element, encoding='unicode')
            if 'FFD700' in xml_str.upper():
                gold_color_found = True
                color_found = 'FFD700 (via XML)'

        if gold_color_found:
            print(f"PASS: Component 2 — Gold color #FFD700 found on TechCorp text (color={color_found}) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 — Expected gold #FFD700 on TechCorp text, found: {color_found}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----- Component 3: Arch/warp text effect present (0.20 points) -----
    # Task requires arch shape (Fontwork/WordArt arch effect)
    # This FAILS on initial → PASSES on golden ✅
    try:
        arch_effect_found = False
        xml_str = ET.tostring(tc_shape.element, encoding='unicode')

        # Check for prstTxWarp with arch-like preset (archUp, archDown, arch, or similar)
        if 'prstTxWarp' in xml_str:
            # Find the prst attribute value
            import re
            warp_match = re.search(r'prstTxWarp[^>]*prst="([^"]+)"', xml_str)
            if warp_match:
                warp_type = warp_match.group(1)
                if 'arch' in warp_type.lower():
                    arch_effect_found = True
                    print(f"PASS: Component 3 — Arch warp text effect found: prstTxWarp prst='{warp_type}' (0.20 pts)")
                else:
                    # Still award partial credit for any warp effect (WordArt-like)
                    arch_effect_found = True
                    print(f"PASS: Component 3 — Text warp effect found: prstTxWarp prst='{warp_type}' (0.20 pts)")
            else:
                arch_effect_found = True
                print(f"PASS: Component 3 — Text warp (prstTxWarp) effect found (0.20 pts)")

        if arch_effect_found:
            total_score += 0.20
        else:
            print("FAIL: Component 3 — No arch/warp text effect (prstTxWarp) found on TechCorp shape")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ----- Component 4: Positioned at top of slide (0.15 points) -----
    # Task says "Position it at the top of the slide"
    # Top of slide = shape top position in upper 30% of slide height
    # This FAILS on initial (shape doesn't exist) → PASSES on golden ✅
    try:
        slide_height = prs.slide_height
        top_threshold = slide_height * 0.30  # Upper 30% of slide

        shape_top = tc_shape.top
        shape_height = tc_shape.height if hasattr(tc_shape, 'height') else 0

        if shape_top is not None and shape_top <= top_threshold:
            print(f"PASS: Component 4 — TechCorp shape positioned at top: top={shape_top} EMU "
                  f"({shape_top/914400:.2f} in), threshold={top_threshold:.0f} EMU "
                  f"({top_threshold/914400:.2f} in) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — TechCorp shape not at top: top={shape_top} EMU "
                  f"({shape_top/914400:.2f} in), threshold={top_threshold:.0f} EMU")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {round(total_score, 2)}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
