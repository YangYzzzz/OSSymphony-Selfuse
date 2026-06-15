"""
Initial Setup: Company intro presentation without Fontwork/WordArt
Task ID: impress_objects_068
Domain: libreoffice_impress
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_objects_068'
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'
DESKTOP_OUTPUT = f'{WORKDIR}/Desktop/company_intro.pptx'


def create_initial():
    prs = Presentation()
    # Set standard widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Company Introduction (dark background, no WordArt) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Dark background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)  # Dark navy blue

    # Subtitle text box (centered, but NO Fontwork/WordArt at top)
    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Innovation for the Future'
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = 'Calibri'
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # Tagline box
    txBox2 = slide1.shapes.add_textbox(Inches(2.0), Inches(4.2), Inches(9.33), Inches(1.0))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = 'Empowering Businesses Worldwide Since 2005'
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = 'Calibri'
    run2.font.size = Pt(20)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)

    # Decorative line shape (rectangle as divider)
    rect = slide1.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(3.5), Inches(4.0), Inches(6.33), Pt(3)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x4A, 0x90, 0xD9)
    rect.line.fill.background()

    # --- Slide 2: About TechCorp ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    # Fill slide 2 background
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    title2 = slide2.shapes.title
    title2.text = 'About TechCorp'
    for para in title2.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.bold = True
            run.font.size = Pt(32)

    content2 = slide2.placeholders[1]
    content2.text = 'Founded in 2005 in Silicon Valley'
    tf_c2 = content2.text_frame
    bullet_items = [
        'Over 5,000 employees across 30 countries',
        'Annual revenue exceeding $2.4 billion',
        'Specializing in enterprise software solutions',
        'ISO 9001 certified since 2010',
        'Awarded Best Tech Employer 2022 & 2023',
    ]
    for item in bullet_items:
        para = tf_c2.add_paragraph()
        para.text = item
        para.level = 1
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xD0, 0xD0, 0xE8)
            run.font.size = Pt(18)

    # --- Slide 3: Products & Services ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Title
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0))
    tf3 = title_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = 'Products & Services'
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.runs[0]
    run3.font.name = 'Calibri'
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Product list
    products = [
        ('TechCorp Cloud Suite', 'Enterprise cloud infrastructure with 99.99% uptime SLA'),
        ('TechCorp Analytics Pro', 'Real-time business intelligence and data visualization'),
        ('TechCorp SecureNet', 'Zero-trust cybersecurity platform for large enterprises'),
        ('TechCorp DevOps Hub', 'Streamlined CI/CD pipeline management and automation'),
    ]
    y_pos = 1.5
    for name, desc in products:
        box = slide3.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(5.5), Inches(0.8))
        tf_p = box.text_frame
        p_name = tf_p.paragraphs[0]
        p_name.text = name
        r_name = p_name.runs[0]
        r_name.font.name = 'Calibri'
        r_name.font.size = Pt(18)
        r_name.font.bold = True
        r_name.font.color.rgb = RGBColor(0x4A, 0x90, 0xD9)

        desc_para = tf_p.add_paragraph()
        desc_para.text = desc
        r_desc = desc_para.runs[0]
        r_desc.font.name = 'Calibri'
        r_desc.font.size = Pt(14)
        r_desc.font.color.rgb = RGBColor(0xC0, 0xC8, 0xE0)
        y_pos += 1.3

    # --- Slide 4: Financial Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0))
    tf4 = title_box4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = 'Financial Highlights 2024'
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.name = 'Calibri'
    run4.font.size = Pt(34)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Financial table
    table_shape = slide4.shapes.add_table(5, 3, Inches(1.5), Inches(1.5), Inches(10.0), Inches(4.5))
    table = table_shape.table

    headers = ['Metric', 'FY 2023', 'FY 2024']
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x6B)

    financial_data = [
        ('Total Revenue', '$2.1B', '$2.4B'),
        ('Net Profit', '$380M', '$445M'),
        ('R&D Investment', '$210M', '$265M'),
        ('New Clients', '1,240', '1,580'),
    ]
    for row_idx, (metric, fy23, fy24) in enumerate(financial_data, 1):
        values = [metric, fy23, fy24]
        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(15)
                    run.font.color.rgb = RGBColor(0xE0, 0xE8, 0xFF)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0x22, 0x28, 0x48)

    # --- Slide 5: Contact ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    contact_box = slide5.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.33), Inches(4.5))
    tf5 = contact_box.text_frame
    tf5.word_wrap = True

    contact_title = tf5.paragraphs[0]
    contact_title.text = 'Contact Us'
    contact_title.alignment = PP_ALIGN.CENTER
    r_ct = contact_title.runs[0]
    r_ct.font.name = 'Calibri'
    r_ct.font.size = Pt(32)
    r_ct.font.bold = True
    r_ct.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    contact_details = [
        'Headquarters: 1 Infinite Loop, Cupertino, CA 95014',
        'Phone: +1 (408) 555-0199',
        'Email: info@techcorp.com',
        'Website: www.techcorp.com',
        'Investor Relations: ir@techcorp.com',
    ]
    for detail in contact_details:
        para = tf5.add_paragraph()
        para.text = detail
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.name = 'Calibri'
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)

    # Save to standard output path
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Also copy to Desktop location as referenced in task context
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    shutil.copy(OUTPUT, DESKTOP_OUTPUT)
    print(f'Copied to Desktop: {DESKTOP_OUTPUT}')


create_initial()
