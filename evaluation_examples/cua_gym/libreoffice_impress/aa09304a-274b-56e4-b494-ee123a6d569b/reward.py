"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm updating my presentation and noticed that slide 4's title needs to be changed to 'Project Timeline'. What's the best way to rename it in LibreOffice Impress?
Generated: 2025-08-07 10:59:40
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_slide4_title(file_path, expected_title):
    print("Starting verification for slide 4 title update...")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File existence (0.2 points)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score}")
        return total_score
    print("✓ File exists (0.2 points)")
    total_score += 0.2

    # Requirement 2: Load presentation (0.2 points)
    try:
        prs = Presentation(file_path)
        print("✓ Presentation loaded (0.2 points)")
        total_score += 0.2
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 3: Verify slide count (>=4 slides) (0.2 points)
    slide_count = len(prs.slides)
    print(f"Slide count: {slide_count}")
    if slide_count >= 4:
        print("✓ Slide count requirement met (0.2 points)")
        total_score += 0.2
    else:
        print("✗ Slide count less than 4")
        print(f"REWARD: {total_score}")
        return total_score

    # Requirement 4: Slide 4 title placeholder exists (0.2 points)
    slide = prs.slides[3]
    title_shape = slide.shapes.title
    if title_shape is None or not hasattr(title_shape, 'text'):
        print("✗ Slide 4 title placeholder not found")
        print(f"REWARD: {total_score}")
        return total_score
    print("✓ Slide 4 title placeholder exists (0.2 points)")
    total_score += 0.2

    # Requirement 5: Title text matches expected (0.2 points)
    actual_title = title_shape.text.strip()
    print(f"Slide 4 title text: '{actual_title}'")
    if actual_title == expected_title:
        print("✓ Slide 4 title matches expected title (0.2 points)")
        total_score += 0.2
    else:
        print(f"✗ Expected title '{expected_title}' but found '{actual_title}'")

    # Final score calculation and output
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/im_updating_my_presentation_and_noticed_that_slide_4s_title_needs_to_be_changed_to_project_timeline_.pptx'
    verify_slide4_title(file_path, 'Project Timeline')
