"""
Reward Script: Apply underline formatting to titles on slides 2, 3, and 4
Task ID: osworld_impress_title_selective_formatting_004
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 2 title is underlined (0.33 pts)
  - Component 2: Slide 3 title is underlined (0.34 pts)
  - Component 3: Slide 4 title is underlined (0.33 pts)
  Total: 1.0 pts
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_004'


def get_title_shape(slide):
    """Return the title placeholder shape from a slide, or None."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            return shape
    return None


def is_title_underlined(slide):
    """
    Check if all non-empty runs in the title placeholder of a slide are underlined.
    Returns True if all runs have underline=True, False otherwise.
    Returns None if no title shape or no runs are found.
    """
    title_shape = get_title_shape(slide)
    if title_shape is None:
        return None

    runs = []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                runs.append(run)

    if not runs:
        return None

    return all(run.font.underline is True for run in runs)


def verify_task(file_path):
    """
    Verify that underline formatting was applied to slides 2, 3, and 4 titles.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: validate slide count
    if len(prs.slides) != 6:
        print(f"FAIL: Expected 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 title is underlined (0.33 points)
    # This FAILS on initial_env (underline=False) and PASSES on golden_env (underline=True)
    try:
        slide2_underlined = is_title_underlined(prs.slides[1])
        if slide2_underlined is True:
            print("PASS: Component 1 — Slide 2 title is underlined (0.33 pts)")
            total_score += 0.33
        elif slide2_underlined is None:
            print("FAIL: Component 1 — Slide 2 title shape not found or has no runs")
        else:
            title_shape = get_title_shape(prs.slides[1])
            title_text = title_shape.text if title_shape else "N/A"
            # Log underline value per run for diagnostics
            runs_info = []
            if title_shape:
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            runs_info.append(f"'{run.text[:30]}'→underline={run.font.underline}")
            print(f"FAIL: Component 1 — Slide 2 title '{title_text}' not underlined. Runs: {runs_info}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 title is underlined (0.34 points)
    # This FAILS on initial_env (underline=False) and PASSES on golden_env (underline=True)
    try:
        slide3_underlined = is_title_underlined(prs.slides[2])
        if slide3_underlined is True:
            print("PASS: Component 2 — Slide 3 title is underlined (0.34 pts)")
            total_score += 0.34
        elif slide3_underlined is None:
            print("FAIL: Component 2 — Slide 3 title shape not found or has no runs")
        else:
            title_shape = get_title_shape(prs.slides[2])
            title_text = title_shape.text if title_shape else "N/A"
            runs_info = []
            if title_shape:
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            runs_info.append(f"'{run.text[:30]}'→underline={run.font.underline}")
            print(f"FAIL: Component 2 — Slide 3 title '{title_text}' not underlined. Runs: {runs_info}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 4 title is underlined (0.33 points)
    # This FAILS on initial_env (underline=False) and PASSES on golden_env (underline=True)
    try:
        slide4_underlined = is_title_underlined(prs.slides[3])
        if slide4_underlined is True:
            print("PASS: Component 3 — Slide 4 title is underlined (0.33 pts)")
            total_score += 0.33
        elif slide4_underlined is None:
            print("FAIL: Component 3 — Slide 4 title shape not found or has no runs")
        else:
            title_shape = get_title_shape(prs.slides[3])
            title_text = title_shape.text if title_shape else "N/A"
            runs_info = []
            if title_shape:
                for para in title_shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            runs_info.append(f"'{run.text[:30]}'→underline={run.font.underline}")
            print(f"FAIL: Component 3 — Slide 4 title '{title_text}' not underlined. Runs: {runs_info}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Informational: verify slides 1, 5, 6 titles are NOT incorrectly underlined
    # (These are pre-conditions; reported only for diagnostics, not for scoring)
    try:
        for slide_idx, slide_num in [(0, 1), (4, 5), (5, 6)]:
            underlined = is_title_underlined(prs.slides[slide_idx])
            if underlined is True:
                print(f"INFO: Slide {slide_num} title is underlined (should remain unchanged per task)")
            else:
                print(f"INFO: Slide {slide_num} title is not underlined (correct — unchanged)")
    except Exception as e:
        print(f"INFO ERROR: Could not check untouched slides: {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
