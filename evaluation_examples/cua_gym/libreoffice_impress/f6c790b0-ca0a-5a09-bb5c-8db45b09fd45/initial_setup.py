"""
Initial Setup: Apply underline formatting to the titles on slides 2, 3, and 4.
Task ID: osworld_impress_title_selective_formatting_004
Domain: libreoffice_impress

Creates a 6-slide conference session deck with Calibri title fonts,
none of which are underlined. The agent must apply underline to slides 2, 3, 4 only.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_run(placeholder, title_text, font_name='Calibri', font_size_pt=36):
    """Set title placeholder text with specified font, no underline."""
    tf = placeholder.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = title_text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = False
    run.font.italic = False
    run.font.underline = False
    return run


def set_body_run(placeholder, body_text, font_name='Calibri', font_size_pt=18):
    """Set body placeholder text with specified font."""
    tf = placeholder.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = body_text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    return run


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    layout_title = prs.slide_layouts[0]  # Title Slide layout
    slide1 = prs.slides.add_slide(layout_title)
    title_ph = slide1.shapes.title
    subtitle_ph = slide1.placeholders[1]

    set_title_run(title_ph, 'Annual Technology Conference 2025', font_size_pt=40)
    set_body_run(subtitle_ph, 'Exploring the Future of Innovation and Connectivity', font_size_pt=22)

    # --- Slide 2: Workshop Overview ---
    layout_content = prs.slide_layouts[1]  # Title and Content layout
    slide2 = prs.slides.add_slide(layout_content)
    title_ph2 = slide2.shapes.title
    content_ph2 = slide2.placeholders[1]

    set_title_run(title_ph2, 'Workshop Overview', font_size_pt=36)
    tf2 = content_ph2.text_frame
    tf2.clear()
    for item in [
        'Three parallel tracks: AI & Machine Learning, Cloud Infrastructure, Cybersecurity',
        'Morning sessions focus on theoretical foundations and industry trends',
        'Afternoon workshops emphasize hands-on practice and case studies',
        'Evening panel discussions with leading industry practitioners',
    ]:
        para = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        run = para.add_run()
        run.text = item
        run.font.name = 'Calibri'
        run.font.size = Pt(18)

    # --- Slide 3: Key Topics and Themes ---
    slide3 = prs.slides.add_slide(layout_content)
    title_ph3 = slide3.shapes.title
    content_ph3 = slide3.placeholders[1]

    set_title_run(title_ph3, 'Key Topics and Themes', font_size_pt=36)
    tf3 = content_ph3.text_frame
    tf3.clear()
    for item in [
        'Generative AI and Large Language Models in enterprise settings',
        'Zero-trust security architectures for distributed teams',
        'Edge computing and real-time data processing pipelines',
        'Sustainable technology practices and green data centers',
        'DevSecOps: integrating security into the development lifecycle',
    ]:
        para = tf3.add_paragraph() if tf3.paragraphs[0].text else tf3.paragraphs[0]
        run = para.add_run()
        run.text = item
        run.font.name = 'Calibri'
        run.font.size = Pt(18)

    # --- Slide 4: Interactive Sessions ---
    slide4 = prs.slides.add_slide(layout_content)
    title_ph4 = slide4.shapes.title
    content_ph4 = slide4.placeholders[1]

    set_title_run(title_ph4, 'Interactive Sessions', font_size_pt=36)
    tf4 = content_ph4.text_frame
    tf4.clear()
    for item in [
        'Live coding challenge: Build a REST API in 60 minutes',
        'Red team vs. blue team cybersecurity exercise',
        'Cloud architecture design sprint with real AWS/Azure scenarios',
        'AI model fine-tuning lab: bring your own dataset',
    ]:
        para = tf4.add_paragraph() if tf4.paragraphs[0].text else tf4.paragraphs[0]
        run = para.add_run()
        run.text = item
        run.font.name = 'Calibri'
        run.font.size = Pt(18)

    # --- Slide 5: Networking Opportunities ---
    slide5 = prs.slides.add_slide(layout_content)
    title_ph5 = slide5.shapes.title
    content_ph5 = slide5.placeholders[1]

    set_title_run(title_ph5, 'Networking Opportunities', font_size_pt=36)
    tf5 = content_ph5.text_frame
    tf5.clear()
    for item in [
        'Speed networking breakfast: meet 20 peers in 40 minutes',
        'Startup showcase: 15 emerging tech companies presenting demos',
        'Recruiter meetup: connect with hiring managers from top firms',
        'Alumni dinner for returning conference attendees',
    ]:
        para = tf5.add_paragraph() if tf5.paragraphs[0].text else tf5.paragraphs[0]
        run = para.add_run()
        run.text = item
        run.font.name = 'Calibri'
        run.font.size = Pt(18)

    # --- Slide 6: Call to Action & Next Steps ---
    slide6 = prs.slides.add_slide(layout_content)
    title_ph6 = slide6.shapes.title
    content_ph6 = slide6.placeholders[1]

    set_title_run(title_ph6, 'Call to Action & Next Steps', font_size_pt=36)
    tf6 = content_ph6.text_frame
    tf6.clear()
    for item in [
        'Register for your preferred workshop tracks by end of week',
        'Submit your session proposal through the conference portal',
        'Join the Slack community for ongoing discussions and resources',
        'Follow up with speakers via LinkedIn using #TechConf2025',
        'Save the date: next conference scheduled for Q2 2026',
    ]:
        para = tf6.add_paragraph() if tf6.paragraphs[0].text else tf6.paragraphs[0]
        run = para.add_run()
        run.text = item
        run.font.name = 'Calibri'
        run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
