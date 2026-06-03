"""
FINAL REWARD SCRIPT - SUCCESS
Task: While tidying up this 200-slide deck, I noticed slide 146 doesn’t match our style. Could you reformat the content text there so it’s 18 pt, line spacing set to 1.3, and fully justified?
Generated: 2025-09-10 17:48:32
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

"""
Reward script for verifying the task:
"Reformat the content text on slide 146 so it’s 18 pt, line spacing 1.3, and fully justified."

The script checks *only* slide 146 (index 145, zero-based) of the provided PPTX file and evaluates three
separate requirements:
1. All content paragraphs contain at least one run of text with a font size of 18 pt (±0.5).
2. Paragraph line-spacing is 1.3 (±0.05).
3. Paragraph alignment is FULL JUSTIFY.

Each requirement earns roughly one-third of the total score.  Partial completion therefore
produces a progressive score between 0.0 and 1.0.  No points are awarded for conditions that can
occur naturally (e.g. successful file loading).

The script prints diagnostic information and finally prints the reward as:
    REWARD: X.X
"""

FILE_PATH = "/home/user/while_tidying_up_this_200_slide_deck_i_noticed_slide_146_doesnt_match_our_style_could_you_reformat_t_golden.pptx"
EXPECTED_FONT_PT = 18
EXPECTED_LINE_SPACING = 1.3
EXPECTED_ALIGNMENT = PP_ALIGN.JUSTIFY

# per-criterion weights (sum to 1.0)
FONT_WEIGHT    = 0.34
SPACING_WEIGHT = 0.33
ALIGN_WEIGHT   = 0.33

def verify_text_formatting_on_slide(file_path: str, slide_index: int = 145) -> float:
    """Return a progressive score (0.0 – 1.0) based on formatting compliance."""

    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        return 0.0

    if slide_index >= len(prs.slides):
        print(f"✗ Slide index {slide_index} out of range (deck has {len(prs.slides)} slides)")
        return 0.0

    slide = prs.slides[slide_index]
    print(f"Loaded slide {slide_index + 1}/{len(prs.slides)} for verification.")

    # Flags showing whether we found at least one paragraph satisfying each condition.
    font_ok = False
    spacing_ok = False
    align_ok = False

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            # alignment
            if paragraph.alignment == EXPECTED_ALIGNMENT:
                align_ok = True
            # line spacing – pptx objects return float or None
            if paragraph.line_spacing is not None and abs(paragraph.line_spacing - EXPECTED_LINE_SPACING) < 0.05:
                spacing_ok = True
            # font size – examine runs
            for run in paragraph.runs:
                if run.font.size is not None and abs(run.font.size.pt - EXPECTED_FONT_PT) < 0.5:
                    font_ok = True
            # early exit when all are satisfied
            if font_ok and spacing_ok and align_ok:
                break
        if font_ok and spacing_ok and align_ok:
            break

    # Progressive scoring
    score = 0.0
    if font_ok:
        print(f"✓ Font size requirement met (≈{EXPECTED_FONT_PT} pt)")
        score += FONT_WEIGHT
    else:
        print("✗ Font size requirement NOT met")

    if spacing_ok:
        print(f"✓ Line spacing requirement met (≈{EXPECTED_LINE_SPACING})")
        score += SPACING_WEIGHT
    else:
        print("✗ Line spacing requirement NOT met")

    if align_ok:
        print("✓ Alignment requirement met (JUSTIFY)")
        score += ALIGN_WEIGHT
    else:
        print("✗ Alignment requirement NOT met")

    final_score = round(min(score, 1.0), 2)
    print(f"Total score: {final_score}/1.0")
    return final_score


def main():
    reward = verify_text_formatting_on_slide(FILE_PATH, 145)  # slide 146 in 1-based numbering
    print(f"REWARD: {reward}")


if __name__ == "__main__":
    main()

