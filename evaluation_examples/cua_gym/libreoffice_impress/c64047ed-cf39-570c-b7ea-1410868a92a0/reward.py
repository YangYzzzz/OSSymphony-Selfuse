"""
FINAL REWARD SCRIPT - SUCCESS
Task: Doing a last-minute brand check in LibreOffice Impress and spotted a rogue header—please set the title on slide 124 to Noto Sans, 42 pt, Bold so it matches the rest.
Generated: 2025-09-10 17:23:47
Status: success
Model: azure-o3
Total Steps: 3
"""

"""Reward script for verifying that the title on slide 124 of the given
PowerPoint file is formatted in Noto Sans, 42 pt, Bold – exactly as
requested in the task instructions.

Scoring (progressive):
  • 0.4 – Every non-empty run in the slide-124 title uses the font
           exactly "Noto Sans" (case-insensitive, trimmed).
  • 0.3 – Font size for every run is ≈ 42 pt (±0.5 pt tolerance to avoid
           floating-point edge cases).
  • 0.3 – Every run is explicitly bold.
  → Perfect compliance yields 1.0.  No points are given for natural
    conditions such as file existence, slide access, etc.

The script prints detailed diagnostic information then prints the final
score in the required format:  "REWARD: X.X".
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

FILE_PATH = (
    "/home/user/doing_a_last_minute_brand_check_in_libreoffice_impress_and_"
    "spotted_a_rogue_headerplease_set_the_titl_golden.pptx"
)

# ---------------------------------------------------------------------
# Helper ----------------------------------------------------------------
# ---------------------------------------------------------------------

def _effective(attr_run, attr_para):
    """Return the run-level attribute if set, else the paragraph-level one."""
    return attr_run if attr_run is not None else attr_para

# ---------------------------------------------------------------------
# Core verification -----------------------------------------------------
# ---------------------------------------------------------------------

def verify_title_format(file_path: str) -> float:
    """Verify slide-124 title formatting and return a progressive score."""

    print(f"Verifying presentation: {file_path}")
    total_score = 0.0  # accumulate up to 1.0

    # ---------- Basic sanity checks (no points awarded) ---------------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to load PPTX: {exc}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 124:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 124 missing")
        print("REWARD: 0.0")
        return 0.0

    # ------------------ Retrieve slide & title ------------------------
    slide = prs.slides[123]  # zero-based index 123 == slide 124

    # Prefer the designated title placeholder; fall back to manual search
    title_shape = slide.shapes.title
    if title_shape is None:
        for sh in slide.shapes:
            if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_shape = sh
                break

    if title_shape is None or not title_shape.has_text_frame:
        print("✗ No title shape with text found on slide 124")
        print("REWARD: 0.0")
        return 0.0

    tf = title_shape.text_frame

    # ----------------- Inspect all text runs --------------------------
    runs_meta = []  # collect dict(name, size_pt, bold) for each run

    for para in tf.paragraphs:
        para_font = para.font  # paragraph-level defaults
        for run in para.runs:
            txt = run.text.strip()
            if not txt:
                continue  # ignore empty runs
            name = _effective(run.font.name, para_font.name)
            size = _effective(run.font.size, para_font.size)
            bold = _effective(run.font.bold, para_font.bold)
            runs_meta.append(
                {
                    "text": txt,
                    "name": name,
                    "size_pt": size.pt if size is not None else None,
                    "bold": bold,
                }
            )

    if not runs_meta:
        print("✗ Title contains no text runs – nothing to verify")
        print("REWARD: 0.0")
        return 0.0

    # --------------- Requirement 1: Font name -------------------------
    font_name_ok = all(
        meta["name"] and meta["name"].strip().lower() == "noto sans" for meta in runs_meta
    )
    if font_name_ok:
        print("✓ All title runs use font ‘Noto Sans’ (0.4)")
        total_score += 0.4
    else:
        print("✗ Font name mismatch in one or more title runs (0)")

    # --------------- Requirement 2: Font size -------------------------
    size_ok = all(
        meta["size_pt"] is not None and abs(meta["size_pt"] - 42) < 0.5 for meta in runs_meta
    )
    if size_ok:
        print("✓ All title runs are ≈ 42 pt (0.3)")
        total_score += 0.3
    else:
        print("✗ Font size mismatch in one or more title runs (0)")

    # --------------- Requirement 3: Bold ------------------------------
    bold_ok = all(meta["bold"] is True for meta in runs_meta)
    if bold_ok:
        print("✓ All title runs are bold (0.3)")
        total_score += 0.3
    else:
        print("✗ Bold property missing in one or more title runs (0)")

    # ------------------------ Wrap-up ---------------------------------
    final = min(total_score, 1.0)
    print(f"Final Score: {final}")
    print(f"REWARD: {final}")
    return final


# ---------------------------------------------------------------------
# Execute verification when script is run ------------------------------
# ---------------------------------------------------------------------
if __name__ == "__main__":
    verify_title_format(FILE_PATH)

