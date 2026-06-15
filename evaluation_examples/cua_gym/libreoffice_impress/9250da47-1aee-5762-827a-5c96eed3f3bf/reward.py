"""
Reward Script: Multiple-choice quiz slide on slide 7 of Bio_Review.pptx
Task ID: impress_teach_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Title 'Quick Check' present on slide 7
  Component 2 (0.35): Quiz question text with all 4 answer options present
  Component 3 (0.20): 'B) Mitochondria' run is bold
  Component 4 (0.20): 'B) Mitochondria' run is colored #2E7D32
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_021'


def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames, including inside groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide7 = prs.slides[6]  # 0-indexed
    text_shapes = get_all_text_shapes(slide7)

    # Collect all text content from slide 7
    all_text_on_slide = ""
    for shape in text_shapes:
        all_text_on_slide += " " + shape.text

    # Component 1: Title 'Quick Check' present on slide 7 (0.25 points)
    # This checks that the title text has been set (it's empty in initial_env)
    try:
        quick_check_found = False
        for shape in text_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'quick check' in run.text.strip().lower():
                            quick_check_found = True
                            break
                    if quick_check_found:
                        break
            if quick_check_found:
                break

        if quick_check_found:
            print(f"PASS: Component 1 - Title 'Quick Check' found on slide 7 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 - Title 'Quick Check' not found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Quiz question and all 4 options present (0.35 points)
    # The initial slide 7 has no quiz text at all - this is entirely task-introduced
    try:
        expected_texts = [
            "what is the powerhouse of the cell",
            "a) nucleus",
            "b) mitochondria",
            "c) ribosome",
            "d) golgi apparatus"
        ]
        all_lower = all_text_on_slide.lower()
        found_count = 0
        for expected in expected_texts:
            if expected in all_lower:
                found_count += 1
                print(f"  Found: '{expected}'")
            else:
                print(f"  Missing: '{expected}'")

        if found_count == len(expected_texts):
            print(f"PASS: Component 2 - All quiz text present ({found_count}/{len(expected_texts)}) (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 2 - Only {found_count}/{len(expected_texts)} text items found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: 'B) Mitochondria' run is bold (0.20 points)
    # In initial_env there is no such text at all, so this will fail on initial
    try:
        b_mito_bold = False
        for shape in text_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and 'b) mitochondria' in run.text.strip().lower():
                            # Check bold - None means inherit (not bold), True means bold
                            if run.font.bold is True:
                                b_mito_bold = True
                                print(f"  Found 'B) Mitochondria' run with bold=True")
                            else:
                                print(f"  Found 'B) Mitochondria' run with bold={run.font.bold}")

        if b_mito_bold:
            print(f"PASS: Component 3 - 'B) Mitochondria' is bold (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 - 'B) Mitochondria' is not bold")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: 'B) Mitochondria' run is colored #2E7D32 (0.20 points)
    # In initial_env there is no such text at all, so this will fail on initial
    try:
        b_mito_green = False
        target_color = "2E7D32"
        for shape in text_shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and 'b) mitochondria' in run.text.strip().lower():
                            try:
                                if run.font.color.type is not None:
                                    actual_color = str(run.font.color.rgb)
                                    if actual_color.upper() == target_color.upper():
                                        b_mito_green = True
                                        print(f"  Found 'B) Mitochondria' with color #{actual_color}")
                                    else:
                                        print(f"  Found 'B) Mitochondria' with wrong color #{actual_color}")
                                else:
                                    print(f"  Found 'B) Mitochondria' but color type is None (no explicit color)")
                            except Exception as ce:
                                print(f"  Color check error: {ce}")

        if b_mito_green:
            print(f"PASS: Component 4 - 'B) Mitochondria' colored #{target_color} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 - 'B) Mitochondria' not colored #{target_color}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
