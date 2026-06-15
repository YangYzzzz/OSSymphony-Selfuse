"""
Initial Setup: Create Bio_Review.pptx with 8 slides, slide 7 blank with title only.
Task ID: impress_teach_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "AP Biology Review"
    slide1.placeholders[1].text = "Comprehensive Unit Review\nMs. Patel - Period 3"

    # --- Slide 2: Cell Structure Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Cell Structure Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Cells are the basic unit of life"
    items2 = [
        "Prokaryotic cells lack a membrane-bound nucleus",
        "Eukaryotic cells contain specialized organelles",
        "All cells are bounded by a plasma membrane",
        "Cell theory was established in the 1830s by Schleiden and Schwann",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 3: Organelle Functions ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Organelle Functions"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Major organelles and their roles:"
    organelles = [
        "Nucleus - stores genetic material (DNA)",
        "Mitochondria - cellular respiration and ATP production",
        "Endoplasmic Reticulum - protein and lipid synthesis",
        "Golgi Apparatus - modifies, packages, and ships proteins",
        "Ribosomes - site of protein synthesis",
        "Lysosomes - intracellular digestion and recycling",
    ]
    for item in organelles:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Cellular Respiration ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Cellular Respiration"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "The process of converting glucose to usable energy"
    steps = [
        "Glycolysis: glucose -> 2 pyruvate (cytoplasm, 2 ATP net)",
        "Krebs Cycle: pyruvate -> CO2 (mitochondrial matrix, 2 ATP)",
        "Electron Transport Chain: NADH/FADH2 -> ~34 ATP (inner membrane)",
        "Overall: C6H12O6 + 6O2 -> 6CO2 + 6H2O + ~38 ATP",
    ]
    for item in steps:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: Photosynthesis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Photosynthesis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Light-dependent and light-independent reactions"
    photo_items = [
        "Light reactions occur in the thylakoid membrane",
        "Water is split to release O2 and produce ATP/NADPH",
        "Calvin Cycle fixes CO2 into glucose (stroma)",
        "Equation: 6CO2 + 6H2O + light -> C6H12O6 + 6O2",
    ]
    for item in photo_items:
        p = body5.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 6: DNA and Protein Synthesis ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "DNA and Protein Synthesis"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "From gene to protein"
    dna_items = [
        "DNA replication occurs during S phase of the cell cycle",
        "Transcription: DNA -> mRNA in the nucleus",
        "Translation: mRNA -> amino acid chain at ribosomes",
        "Codons are 3-nucleotide sequences specifying amino acids",
    ]
    for item in dna_items:
        p = body6.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: BLANK with title placeholder only ---
    # Use Title Only layout (index 5 = Blank; we use 5 for truly blank)
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add a title text box manually since blank layout has no placeholders
    title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = ""
    # Leave it empty - the task is to fill in the title and add quiz content

    # --- Slide 8: Study Tips ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Study Tips for the Exam"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Prepare effectively for your biology exam"
    tips = [
        "Review all vocabulary terms from each unit",
        "Practice drawing and labeling cell diagrams",
        "Work through past exam questions",
        "Form study groups to discuss complex topics",
        "Focus on understanding processes, not just memorizing facts",
    ]
    for item in tips:
        p = body8.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
