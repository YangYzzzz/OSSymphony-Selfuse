"""
FINAL REWARD SCRIPT - SUCCESS
Task: Change the first two paragraphs to 1.5 line spacing.
Generated: 2025-10-17 12:13:59
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

# -------------------------------------------------------------
# Reward Script: Verify that ONLY the first two paragraphs in the
# presentation have been changed to 1.5-line spacing
# -------------------------------------------------------------
# Scoring rules (progressive, max = 1.0):
#   • 0.4 pts  – first paragraph has 1.5 line spacing
#   • 0.4 pts  – second paragraph has 1.5 line spacing
#   • 0.2 pts  – no OTHER paragraph in the file was set to 1.5
#                 (i.e. change limited strictly to the first two)
# NOTE: No points are given for file existence or successful load –
#       those are prerequisites, not achievements.
# -------------------------------------------------------------

FILE_PATH = "/home/user/change_the_first_two_paragraphs_to_15_line_spacing.pptx"

# Helper: decide if a spacing value represents "1.5 lines"

def is_spacing_one_point_five(spacing):
    """Return True if spacing corresponds to 1.5-line spacing.
    python-pptx stores line_spacing as a float multiple (e.g., 1.5) or
    sometimes as an integer EMU value ~150000. Accept a small tolerance.
    """
    if spacing is None:
        return False
    # Float multiple case
    if isinstance(spacing, (float, int)):
        # Accept floats 1.45–1.55 OR raw ints around 150000
        if 1.45 <= float(spacing) <= 1.55:
            return True
        if 149000 <= float(spacing) <= 151000:
            return True
    return False

# Gather all non-empty paragraphs in the presentation (in order)

def collect_paragraphs(pres):
    paragraphs = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                # Build text string from runs
                text = "".join(run.text for run in para.runs).strip()
                if text:  # keep only non-empty paragraphs
                    paragraphs.append((para, text))
    return paragraphs

# Core verification routine

def verify_line_spacing(file_path):
    total_score = 0.0
    max_score = 1.0

    # 1. Ensure file exists & is loadable (prerequisite, no points)
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation ({len(prs.slides)} slides)")
    except Exception as e:
        print("✗ Could not load presentation:", e)
        return 0.0

    # 2. Collect paragraphs
    paragraphs = collect_paragraphs(prs)
    print(f"Total non-empty paragraphs detected: {len(paragraphs)}")

    if len(paragraphs) < 2:
        print("✗ Need at least two paragraphs to evaluate task")
        return 0.0

    first_two = paragraphs[:2]
    remaining = paragraphs[2:]

    # 3. Check first two paragraphs
    for idx, (para, text) in enumerate(first_two, start=1):
        spacing_ok = is_spacing_one_point_five(para.line_spacing)
        print(
            f"Paragraph {idx}: '{text[:30]}...'  spacing={para.line_spacing} -> "
            f"{'OK' if spacing_ok else 'WRONG'}"
        )
        if spacing_ok:
            total_score += 0.4

    # 4. Confirm no other paragraph was (incorrectly) set to 1.5
    others_ok = True
    for para, text in remaining:
        if is_spacing_one_point_five(para.line_spacing):
            others_ok = False
            print(
                f"✗ Additional paragraph has 1.5 spacing (should remain unchanged): "
                f"'{text[:30]}...'"
            )
            break

    if others_ok:
        if remaining:
            print("✓ All other paragraphs retain original spacing")
        else:
            print("✓ Only two paragraphs present and both verified")
        total_score += 0.2

    # 5. Cap score at 1.0 and report
    final = min(total_score, max_score)
    print(f"Total score: {final}")
    return final


if __name__ == "__main__":
    reward = verify_line_spacing(FILE_PATH)
    print(f"REWARD: {reward}")

