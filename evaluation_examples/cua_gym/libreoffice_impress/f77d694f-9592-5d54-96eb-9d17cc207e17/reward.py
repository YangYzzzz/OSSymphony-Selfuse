"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set every italic word in paragraphs 2–4 to size 13 pt.
Generated: 2025-10-17 13:23:23
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.util import Pt
import os


def verify_italic_word_size(file_path: str, target_size_pt: int = 13) -> float:
    """Verify that every italic word (run) in paragraphs 2–4 of every text frame
    across the presentation is set to the target font size.

    Scoring rules:
        - 1.0  : every italic run in paragraphs 2–4 is exactly target_size_pt
        - <1.0 : proportional to the fraction of correctly-sized italic runs
        - 0.0  : no runs are correct or file can’t be loaded
    """

    print(f"Verifying italic words size in presentation: {file_path}")

    # ------------------------------------------------------------------
    # 1. Basic file checks (NO points awarded here – prerequisite only)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Iterate through slides → shapes → text frames → paragraphs 2–4
    # ------------------------------------------------------------------
    total_italic_runs = 0
    correct_italic_runs = 0
    target_size = Pt(target_size_pt)

    def text_frames_iterator(shp):
        """Yield all text frames contained in *shp* (shape or table)."""
        if hasattr(shp, "has_text_frame") and shp.has_text_frame:
            yield shp.text_frame
        # Tables (shape_type 19)
        if getattr(shp, "shape_type", None) == 19:  # MSO_SHAPE_TYPE.TABLE
            for row in shp.table.rows:
                for cell in row.cells:
                    yield cell.text_frame

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            for tf in text_frames_iterator(shape):
                paragraphs = tf.paragraphs
                # paragraphs 2-4 correspond to indices 1-3
                for p_idx in range(1, 4):
                    if p_idx >= len(paragraphs):
                        continue  # paragraph doesn’t exist
                    paragraph = paragraphs[p_idx]
                    for run in paragraph.runs:
                        if run.font.italic is True:  # only explicitly italic runs
                            total_italic_runs += 1
                            if run.font.size == target_size:
                                correct_italic_runs += 1
                            else:
                                print(
                                    f"✗ Slide {slide_idx}, paragraph {p_idx + 1}: "
                                    f"Run '{run.text}' size {run.font.size} ≠ {target_size}"
                                )

    # ------------------------------------------------------------------
    # 3. Scoring
    # ------------------------------------------------------------------
    if total_italic_runs == 0:
        # No italic words in the specified paragraphs → requirement trivially met
        print("No italic runs found in paragraphs 2-4; treating as fully correct.")
        return 1.0

    score = correct_italic_runs / total_italic_runs

    print(f"Total italic runs in paragraphs 2-4: {total_italic_runs}")
    print(f"Runs with correct size {target_size_pt} pt: {correct_italic_runs}")
    print(f"Computed score: {score}")
    return score


if __name__ == "__main__":
    FILE_PATH = "/home/user/set_every_italic_word_in_paragraphs_24_to_size_13_pt.pptx"
    reward = verify_italic_word_size(FILE_PATH)
    print(f"REWARD: {reward}")
