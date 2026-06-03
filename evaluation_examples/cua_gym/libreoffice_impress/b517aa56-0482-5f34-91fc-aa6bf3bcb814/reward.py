"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 113 looks odd because the heading is still black. Can you show me how to use LibreOffice Impress’s eyedropper to copy the exact #FF5733 fill color from Rectangle 1 and apply it to the title text?
Generated: 2025-09-10 23:05:56
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor

"""
Reward Script for LibreOffice Impress Eyedropper Task
----------------------------------------------------
Task Recap:
Slide 113 should have its title text color changed from black to the exact fill color
(#FF5733) used by “Rectangle 1”.  The reward script verifies:
1.  “Rectangle 1” exists on slide 113 and is filled with #FF5733.
2.  The title text on slide 113 has been recolored to #FF5733.
Progressive scoring:
• 0.4 pts – Rectangle 1 with correct fill color
• 0.6 pts – Title text recolored correctly
Total possible = 1.0 (perfect completion)
The script prints detailed diagnostics and always outputs
"REWARD: X.X" where X.X ∈ [0.0, 1.0].
"""

def verify_task(file_path: str) -> float:
    max_score = 1.0
    score = 0.0
    target_hex = "FF5733"  # Expected colour

    # ---------- Load presentation ----------
    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # ---------- Locate slide 113 ----------
    target_index = 112  # 0-based index for slide 113
    if len(prs.slides) <= target_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides; slide 113 missing")
        return 0.0
    slide = prs.slides[target_index]

    # ---------- Requirement 1: Rectangle 1 fill colour ----------
    rectangle_found = False
    rectangle_correct_colour = False
    for shape in slide.shapes:
        name = getattr(shape, "name", "").strip()
        if name.lower() == "rectangle 1":
            rectangle_found = True
            try:
                if shape.fill and shape.fill.type == 1:  # Solid fill
                    rgb = shape.fill.fore_color.rgb
                    print(f"Rectangle 1 fill rgb: {rgb}")
                    if rgb == RGBColor.from_string(target_hex):
                        rectangle_correct_colour = True
            except Exception as e:
                print(f"Error checking rectangle fill: {e}")
            break

    if rectangle_found and rectangle_correct_colour:
        print("✓ Rectangle 1 has correct fill colour (#FF5733) (0.4 pts)")
        score += 0.4
    else:
        if not rectangle_found:
            print("✗ Rectangle 1 not found on slide 113")
        elif not rectangle_correct_colour:
            print("✗ Rectangle 1 fill colour is incorrect")

    # ---------- Requirement 2: Title text colour ----------
    title_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            title_shape = shape
            break
    if title_shape is None:  # Fallback – first text shape as probable title
        for shape in slide.shapes:
            if getattr(shape, "text_frame", None) and shape.text_frame.text.strip():
                title_shape = shape
                break

    title_correct_colour = False
    if title_shape and getattr(title_shape, "text_frame", None):
        for paragraph in title_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                rgb = run.font.color.rgb if run.font and run.font.color else None
                if rgb:
                    print(f"Title run colour: {rgb}")
                    if rgb == RGBColor.from_string(target_hex):
                        title_correct_colour = True
                        break
            if title_correct_colour:
                break

    if title_correct_colour:
        print("✓ Title text colour correctly set to #FF5733 (0.6 pts)")
        score += 0.6
    else:
        print("✗ Title text colour not set to #FF5733")

    # ---------- Final score ----------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_113_looks_odd_because_the_heading_is_still_black_can_you_show_me_how_to_use_libreoffice_impres_golden.pptx"
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")

