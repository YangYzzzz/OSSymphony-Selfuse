"""
Reward Script: Set up two-column layout on slide 4 with colored content areas
Task ID: impress_el_084
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Slide 4 uses 'Two Content' layout
  Component 2 (0.3): Left content placeholder (idx=1) has solid fill #E3F2FD
  Component 3 (0.3): Right content placeholder (idx=2) has solid fill #E8F5E9
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_el_084'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4, 0-indexed

    # Component 1: Slide 4 uses 'Two Content' layout (0.4 points)
    try:
        layout_name = slide.slide_layout.name if slide.slide_layout else None
        if layout_name == 'Two Content':
            print(f"PASS: Component 1 — Slide 4 layout is 'Two Content' (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Expected layout 'Two Content', found '{layout_name}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Build a map of placeholder idx -> shape for slide 4
    placeholders = {}
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholders[shape.placeholder_format.idx] = shape

    # Component 2: Left content placeholder (idx=1) has solid fill #E3F2FD (0.3 points)
    try:
        if 1 in placeholders:
            shape = placeholders[1]
            fill = shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID
                rgb = str(fill.fore_color.rgb)
                if rgb.upper() == 'E3F2FD':
                    print(f"PASS: Component 2 — Left content placeholder fill is #E3F2FD (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 2 — Left content placeholder fill is #{rgb}, expected #E3F2FD")
            else:
                print(f"FAIL: Component 2 — Left content placeholder fill type is {fill.type}, expected SOLID (1)")
        else:
            print(f"FAIL: Component 2 — No placeholder with idx=1 found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Right content placeholder (idx=2) has solid fill #E8F5E9 (0.3 points)
    try:
        if 2 in placeholders:
            shape = placeholders[2]
            fill = shape.fill
            if fill.type is not None and fill.type == 1:  # SOLID
                rgb = str(fill.fore_color.rgb)
                if rgb.upper() == 'E8F5E9':
                    print(f"PASS: Component 3 — Right content placeholder fill is #E8F5E9 (0.3 pts)")
                    total_score += 0.3
                else:
                    print(f"FAIL: Component 3 — Right content placeholder fill is #{rgb}, expected #E8F5E9")
            else:
                print(f"FAIL: Component 3 — Right content placeholder fill type is {fill.type}, expected SOLID (1)")
        else:
            print(f"FAIL: Component 3 — No placeholder with idx=2 found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
