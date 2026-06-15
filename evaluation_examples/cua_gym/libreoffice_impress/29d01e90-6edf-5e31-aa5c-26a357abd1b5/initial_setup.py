"""
Initial Setup: Create Comparison_Deck presentation with 8 slides, slide 4 using Blank layout.
Task ID: impress_el_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, text, font_size=18, bold=False, color=None):
    """Helper to set text on a placeholder with formatting."""
    placeholder.text = text
    for para in placeholder.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor(*color)


def add_bullet_content(text_frame, items, font_size=14):
    """Add bullet items to a text frame."""
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide1.shapes.title.text = "Comparison Deck"
    slide1.placeholders[1].text = "Q4 2025 Product Strategy Review"

    # --- Slide 2: Title and Content - Market Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "Market Overview"
    body2 = slide2.placeholders[1].text_frame
    add_bullet_content(body2, [
        "Global market size reached $48.2B in 2025",
        "Year-over-year growth of 12.3%",
        "North America leads with 38% market share",
        "APAC region showing fastest growth at 18.7%",
        "Key competitors: TechVance, NovaSoft, CloudPeak",
    ])

    # --- Slide 3: Title and Content - Product Roadmap ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide3.shapes.title.text = "Product Roadmap"
    body3 = slide3.placeholders[1].text_frame
    add_bullet_content(body3, [
        "Phase 1: Core platform redesign (Jan-Mar 2026)",
        "Phase 2: AI integration module (Apr-Jun 2026)",
        "Phase 3: Enterprise SSO & compliance (Jul-Sep 2026)",
        "Phase 4: Analytics dashboard v3.0 (Oct-Dec 2026)",
        "Stretch goal: Mobile-first experience refresh",
    ])

    # --- Slide 4: Blank layout (THIS IS THE KEY SLIDE) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Add a simple textbox with content to make it non-empty
    txBox = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Feature Comparison: Plan A vs Plan B"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(28)
        run.font.bold = True

    # Add some placeholder text content
    txBox2 = slide4.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(4), Inches(3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_bullet_content(tf2, [
        "Plan A: Standard Tier",
        "Up to 50 users",
        "Basic analytics",
        "Email support only",
    ])

    txBox3 = slide4.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4), Inches(3))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    add_bullet_content(tf3, [
        "Plan B: Enterprise Tier",
        "Unlimited users",
        "Advanced analytics + AI",
        "24/7 priority support",
    ])

    # --- Slide 5: Title and Content - Revenue Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide5.shapes.title.text = "Revenue Analysis"
    body5 = slide5.placeholders[1].text_frame
    add_bullet_content(body5, [
        "Total revenue: $12.8M (up 15% YoY)",
        "Recurring revenue: $9.4M (73% of total)",
        "New customer acquisition: 342 accounts",
        "Average deal size increased to $37.4K",
        "Churn rate decreased to 4.2%",
    ])

    # --- Slide 6: Title Only - Team Structure ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide6.shapes.title.text = "Team Structure"
    txBox6 = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(4))
    tf6 = txBox6.text_frame
    tf6.word_wrap = True
    add_bullet_content(tf6, [
        "Engineering: 45 engineers across 6 squads",
        "Product: 8 product managers, 4 designers",
        "Sales: 22 account executives, 5 SDRs",
        "Customer Success: 12 CSMs, 3 support leads",
        "Data Science: 6 ML engineers, 2 analysts",
    ])

    # --- Slide 7: Title and Content - Key Risks ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide7.shapes.title.text = "Key Risks & Mitigations"
    body7 = slide7.placeholders[1].text_frame
    add_bullet_content(body7, [
        "Risk: Competitor launch of similar AI module",
        "Mitigation: Accelerate Phase 2 delivery by 3 weeks",
        "Risk: Key talent attrition in engineering",
        "Mitigation: Retention bonuses and mentorship program",
        "Risk: Regulatory changes in APAC data residency",
    ])

    # --- Slide 8: Title Slide - Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
