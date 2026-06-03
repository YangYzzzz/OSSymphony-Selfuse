"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m at slide 106 and need to start a fresh section. How do I drop in a new slide right after it and make the title pick up the exact same formatting that’s on the main heading in slide 1?
Generated: 2025-09-10 16:14:45
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
import os
from pptx.dml.color import RGBColor

# ---------------------------------------------------------------
# Reward Script for:
# "I’m at slide 106 and need to start a fresh section. How do I drop
#  in a new slide right after it and make the title pick up the exact
#  same formatting that’s on the main heading in slide 1?"
# ---------------------------------------------------------------
# Verification Logic
#   1. A NEW slide (making total slides >= 107) must exist.
#   2. The title on that new slide (slide index 106) must match the
#      formatting of the main heading on slide 1 (index 0).
#   3. Font attributes checked: size, bold, font-name, colour.
#   4. Progressive scoring awarded per attribute plus slide-addition.
#      Perfect match -> REWARD 1.0
# ---------------------------------------------------------------

def _extract_first_run_font(shape):
    """Return a dict of font attributes (name, size, bold, color) from the
    first run of the first paragraph in a text shape.
    """
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return None

    tf = shape.text_frame
    if not tf.paragraphs:
        return None

    p = tf.paragraphs[0]
    if not p.runs:
        return None

    run = p.runs[0]
    font = run.font

    attrs = {
        "name": (font.name or "").lower() if font.name else None,
        "size": font.size.pt if font.size else None,
        "bold": bool(font.bold) if font.bold is not None else False,
        "color": tuple(font.color.rgb) if font.color and font.color.rgb else None,
    }
    return attrs


def _slide_first_text_attrs(slide):
    """Find the first text-containing shape on the slide and return its
    font attributes.
    """
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            attrs = _extract_first_run_font(shape)
            if attrs:
                return attrs
    return None


def _compare_attrs(a1, a2):
    """Compare two font-attribute dicts. Returns a dict mapping each checked
    attribute to a boolean indicating match (size tolerance ±1pt)."""
    if a1 is None or a2 is None:
        return {k: False for k in ["size", "bold", "name", "color"]}

    return {
        "size": a1["size"] is not None and a2["size"] is not None and abs(a1["size"] - a2["size"]) <= 1,
        "bold": a1["bold"] == a2["bold"],
        "name": a1["name"] == a2["name"],
        "color": a1["color"] == a2["color"],
    }


def verify_task(file_path: str) -> float:
    """Main verification function. Returns a float reward 0.0 – 1.0."""
    print(f"Verifying task on file: {file_path}")
    score = 0.0

    # ----- Load presentation -------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Presentation loaded. Slide count: {total_slides}")

    # ----- Requirement 1: New slide added ------------------------------------
    if total_slides >= 107:
        print("✓ Slide count indicates a new slide was added (>=107)")
        score += 0.4  # 40% of total score
    else:
        print("✗ Slide count < 107 — new slide not detected")

    # ----- Requirement 2: Formatting match -----------------------------------
    # Get attributes from slide 1 heading
    ref_attrs = _slide_first_text_attrs(prs.slides[0])
    if not ref_attrs:
        print("✗ Could not extract heading formatting from slide 1")
        print(f"REWARD: {score}")
        return score
    print("Slide 1 heading attrs:", ref_attrs)

    # Ensure slide 107 exists
    target_index = 106  # zero-based index
    if total_slides <= target_index:
        print(f"✗ Slide {target_index+1} (new section slide) not found")
        print(f"REWARD: {score}")
        return score

    new_attrs = _slide_first_text_attrs(prs.slides[target_index])
    if not new_attrs:
        print("✗ Could not extract title formatting from new slide")
        print(f"REWARD: {score}")
        return score
    print(f"Slide {target_index+1} title attrs:", new_attrs)

    match = _compare_attrs(ref_attrs, new_attrs)
    print("Attribute match results:", match)

    # Weighted progressive scoring for each matched attribute
    weights = {"size": 0.2, "bold": 0.15, "name": 0.10, "color": 0.15}
    for attr, ok in match.items():
        if ok:
            score += weights[attr]
            print(f"✓ {attr} matches (+{weights[attr]:.2f})")
        else:
            print(f"✗ {attr} does NOT match (+0.00)")

    # ----- Final score -------------------------------------------------------
    final = min(score, 1.0)
    print(f"Total score: {final:.2f}")
    print(f"REWARD: {final}")
    return final


if __name__ == "__main__":
    FILE = "/home/user/im_at_slide_106_and_need_to_start_a_fresh_section_how_do_i_drop_in_a_new_slide_right_after_it_and_ma_golden.pptx"
    verify_task(FILE)

