"""
Initial Setup: Travel Agency presentation with airplane icon on slide 2
Task ID: impress_ma_063
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
AIRPLANE_IMG = f'{WORKDIR}/airplane_icon.png'


def create_airplane_icon():
    """Create a simple airplane icon PNG using Pillow."""
    img = Image.new('RGBA', (200, 200), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    # Airplane body
    draw.polygon([
        (100, 20),   # nose
        (90, 60),
        (40, 100),   # left wing tip
        (80, 95),
        (70, 140),
        (50, 170),   # left tail
        (85, 150),
        (100, 180),  # tail bottom
        (115, 150),
        (150, 170),  # right tail
        (130, 140),
        (120, 95),
        (160, 100),  # right wing tip
        (110, 60),
    ], fill=(41, 128, 185, 255), outline=(30, 90, 140, 255))

    # Window dots
    for y in [55, 70, 85]:
        draw.ellipse([96, y, 104, y + 5], fill=(255, 255, 255, 220))

    img.save(AIRPLANE_IMG)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    create_airplane_icon()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Wanderlust Travel Agency"
    slide1.placeholders[1].text = "Your Journey Begins Here"
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0x00, 0x3F, 0x72)
    # Make title text white
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(44)
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        run.font.size = Pt(24)

    # --- Slide 2: Popular Destinations (with airplane icon) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Popular Destinations 2025"
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0x00, 0x3F, 0x72)

    # Destination cards as text boxes
    destinations = [
        ("Santorini, Greece", "$2,450 per person\n7 nights at Caldera View Resort\nIncludes flights & transfers"),
        ("Kyoto, Japan", "$3,100 per person\n10 nights at Sakura Garden Inn\nIncludes rail pass & guided tours"),
        ("Machu Picchu, Peru", "$2,800 per person\n8 nights adventure package\nIncludes Inca Trail permit"),
    ]
    for i, (title, details) in enumerate(destinations):
        x = Inches(0.5 + i * 4.0)
        y = Inches(1.5)
        box = slide2.shapes.add_textbox(x, y, Inches(3.5), Inches(3.0))
        tf = box.text_frame
        tf.word_wrap = True
        p_title = tf.paragraphs[0]
        p_title.text = title
        r_title = p_title.runs[0]
        r_title.font.size = Pt(20)
        r_title.font.bold = True
        r_title.font.color.rgb = RGBColor(0x1A, 0x5E, 0x8F)
        p_det = tf.add_paragraph()
        p_det.text = details
        for r in p_det.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Airplane icon at bottom-left (x=1in, y=6in) - this is the target shape
    airplane_pic = slide2.shapes.add_picture(
        AIRPLANE_IMG,
        Inches(1), Inches(6),
        Inches(0.8), Inches(0.8)
    )
    airplane_pic.name = "Airplane Icon"

    # --- Slide 3: Our Services ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Our Services"
    run3 = p3.runs[0]
    run3.font.size = Pt(36)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x00, 0x3F, 0x72)

    services = [
        "Custom Itinerary Planning - Tailored trips designed around your preferences",
        "Group Travel Coordination - Corporate retreats, family reunions, weddings",
        "Luxury Concierge Services - Private transfers, VIP experiences, fine dining",
        "Travel Insurance Packages - Comprehensive coverage for peace of mind",
        "Visa & Documentation Support - Hassle-free travel paperwork assistance",
        "24/7 Emergency Support - Always available when you need us most",
    ]
    for j, svc in enumerate(services):
        sbox = slide3.shapes.add_textbox(Inches(1), Inches(1.3 + j * 0.9), Inches(10), Inches(0.8))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = svc
        for r in sp.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # --- Slide 4: Summer 2025 Deals ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Summer 2025 Special Deals"
    run4 = p4.runs[0]
    run4.font.size = Pt(36)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)

    # Table of deals
    table_shape = slide4.shapes.add_table(5, 4, Inches(0.8), Inches(1.5), Inches(11), Inches(4))
    table = table_shape.table
    headers = ["Package", "Duration", "Price", "Savings"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for r in cell.text_frame.paragraphs[0].runs:
            r.font.bold = True
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xC0, 0x39, 0x2B)

    deals = [
        ["Mediterranean Cruise", "14 nights", "$4,299", "Save 25%"],
        ["Bali Beach Retreat", "10 nights", "$2,199", "Save 30%"],
        ["African Safari Adventure", "12 nights", "$5,450", "Save 20%"],
        ["Iceland Northern Lights", "7 nights", "$3,150", "Save 15%"],
    ]
    for r_idx, row_data in enumerate(deals, 1):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            for r in cell.text_frame.paragraphs[0].runs:
                r.font.size = Pt(14)

    # --- Slide 5: Client Testimonials ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "What Our Clients Say"
    run5 = p5.runs[0]
    run5.font.size = Pt(36)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x00, 0x3F, 0x72)

    testimonials = [
        ('"Wanderlust planned our honeymoon to the Maldives and it was absolutely magical. Every detail was perfect."',
         "— Elena & Marco Rodriguez, San Francisco"),
        ('"Our corporate retreat in Tuscany was a huge success. The team building activities and logistics were handled flawlessly."',
         "— James Whitfield, VP Operations, TechNova Inc."),
        ('"I\'ve used Wanderlust for 5 family vacations. Their attention to kid-friendly activities and dietary needs is unmatched."',
         "— Priya Sharma, Chicago"),
    ]
    for k, (quote, author) in enumerate(testimonials):
        tbox = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5 + k * 1.8), Inches(10), Inches(1.5))
        ttf = tbox.text_frame
        ttf.word_wrap = True
        pq = ttf.paragraphs[0]
        pq.text = quote
        for r in pq.runs:
            r.font.size = Pt(16)
            r.font.italic = True
            r.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        pa = ttf.add_paragraph()
        pa.text = author
        pa.alignment = PP_ALIGN.RIGHT
        for r in pa.runs:
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x1A, 0x5E, 0x8F)

    # --- Slide 6: Our Team ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Meet Our Expert Team"
    run6 = p6.runs[0]
    run6.font.size = Pt(36)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x00, 0x3F, 0x72)

    team = [
        ("Sophia Nakamura", "Founder & CEO", "20+ years in luxury travel, former Condé Nast editor"),
        ("David Okonkwo", "Head of Operations", "Specializes in African & Middle East destinations"),
        ("Isabella Torres", "Lead Consultant", "Expert in European cultural & culinary tours"),
        ("Raj Patel", "Technology Director", "Building seamless digital booking experiences"),
    ]
    for m, (name, role, desc) in enumerate(team):
        x = Inches(0.5 + (m % 2) * 6.2)
        y = Inches(1.5 + (m // 2) * 2.8)
        mbox = slide6.shapes.add_textbox(x, y, Inches(5.5), Inches(2.2))
        mtf = mbox.text_frame
        mtf.word_wrap = True
        pn = mtf.paragraphs[0]
        pn.text = name
        for r in pn.runs:
            r.font.size = Pt(22)
            r.font.bold = True
            r.font.color.rgb = RGBColor(0x00, 0x3F, 0x72)
        pr = mtf.add_paragraph()
        pr.text = role
        for r in pr.runs:
            r.font.size = Pt(16)
            r.font.italic = True
            r.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
        pd = mtf.add_paragraph()
        pd.text = desc
        for r in pd.runs:
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 7: Contact Us ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    bg7 = slide7.background.fill
    bg7.solid()
    bg7.fore_color.rgb = RGBColor(0x00, 0x3F, 0x72)

    txBox7 = slide7.shapes.add_textbox(Inches(2), Inches(1), Inches(9), Inches(1.5))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Start Your Journey Today"
    p7.alignment = PP_ALIGN.CENTER
    run7 = p7.runs[0]
    run7.font.size = Pt(40)
    run7.font.bold = True
    run7.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    contact_info = [
        "Phone: +1 (555) 982-4567",
        "Email: bookings@wanderlusttravel.com",
        "Web: www.wanderlusttravel.com",
        "Office: 742 Pacific Avenue, Suite 300, San Francisco, CA 94115",
    ]
    for ci in contact_info:
        cbox = slide7.shapes.add_textbox(Inches(3), Inches(2.8 + contact_info.index(ci) * 0.8), Inches(7), Inches(0.7))
        ctf = cbox.text_frame
        cp = ctf.paragraphs[0]
        cp.text = ci
        cp.alignment = PP_ALIGN.CENTER
        for r in cp.runs:
            r.font.size = Pt(20)
            r.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
