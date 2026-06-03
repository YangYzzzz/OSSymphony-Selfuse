"""
Reward Script: Multi-slide build animation for system architecture
Task ID: impress_stu_090
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide 3 - Client Layer content (title + 3 boxes)
  Component 2 (0.20): Slide 4 - Layers 1 & 2 content (title + client + server boxes + arrows)
  Component 3 (0.20): Slide 5 - Complete Architecture content (title + all layers + arrows)
  Component 4 (0.15): Color differentiation per layer (blue/green/orange)
  Component 5 (0.10): Connecting arrows between layers
  Component 6 (0.15): Morph transitions on slides 4 and 5
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_090'


def classify_color(rgb_str):
    """Classify an RGB hex string into a layer color category."""
    if rgb_str is None:
        return None
    rgb_str = rgb_str.upper()
    r = int(rgb_str[0:2], 16)
    g = int(rgb_str[2:4], 16)
    b = int(rgb_str[4:6], 16)
    # Blue family: high blue, lower red/green
    if b > 150 and b > r and b > g:
        return 'blue'
    # Green family: high green, lower red/blue
    if g > 100 and g > r and g > b:
        return 'green'
    # Orange family: high red, medium green, low blue
    if r > 180 and g > 80 and b < 100:
        return 'orange'
    return None


def get_shape_fill_color(shape):
    """Get the fill color of a shape as hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def check_morph_transition(pptx_path, slide_num):
    """Check if a slide (1-indexed) has a Morph transition."""
    morph_ns = 'http://schemas.microsoft.com/office/powerpoint/2015/main'
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            with zf.open(f'ppt/slides/slide{slide_num}.xml') as f:
                root = ET.parse(f).getroot()
                # Look for transition element
                for elem in root.iter():
                    if 'transition' in elem.tag:
                        # Look for morph child in any namespace
                        for child in elem:
                            if 'morph' in child.tag.lower():
                                return True
        except (KeyError, ET.ParseError):
            pass
    return False


def get_autoshapes_with_text(slide):
    """Get all AUTO_SHAPE shapes from a slide, returning (text, fill_color) tuples."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    results = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = ''
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            fill_color = get_shape_fill_color(shape)
            results.append((text, fill_color))
    return results


def get_slide_title(slide):
    """Get the title text from the first text box on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ''


def verify_task(file_path):
    """Verify task completion with progressive scoring."""
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]
    slide4 = prs.slides[3]
    slide5 = prs.slides[4]

    # --- Component 1: Slide 3 - Client Layer (0.20 points) ---
    # Must have title containing "Layer 1" and 3 client boxes (Web, Mobile, API)
    try:
        title3 = get_slide_title(slide3)
        has_title3 = 'layer 1' in title3.lower() or 'architecture' in title3.lower()

        autoshapes3 = get_autoshapes_with_text(slide3)
        texts3 = {t.lower() for t, _ in autoshapes3 if t}
        client_boxes = {'web', 'mobile', 'api'}
        has_client = client_boxes.issubset(texts3)

        if has_title3 and has_client:
            print(f"PASS: Component 1 - Slide 3 has title '{title3}' and client boxes {texts3} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 - Slide 3: title='{title3}' (ok={has_title3}), "
                  f"client boxes found={texts3 & client_boxes} (ok={has_client})")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # --- Component 2: Slide 4 - Layers 1 & 2 (0.20 points) ---
    # Must have title, client boxes, and server boxes (Auth, Logic, Cache)
    try:
        title4 = get_slide_title(slide4)
        has_title4 = ('layer' in title4.lower() and '2' in title4) or 'architecture' in title4.lower()

        autoshapes4 = get_autoshapes_with_text(slide4)
        texts4 = {t.lower() for t, _ in autoshapes4 if t}
        server_boxes = {'auth', 'logic', 'cache'}
        has_client4 = client_boxes.issubset(texts4)
        has_server4 = server_boxes.issubset(texts4)

        if has_title4 and has_client4 and has_server4:
            print(f"PASS: Component 2 - Slide 4 has title '{title4}', client and server boxes (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 - Slide 4: title='{title4}' (ok={has_title4}), "
                  f"client={has_client4}, server={has_server4}, texts={texts4}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # --- Component 3: Slide 5 - Complete Architecture (0.20 points) ---
    # Must have title, all three layers' boxes
    try:
        title5 = get_slide_title(slide5)
        has_title5 = 'complete' in title5.lower() or 'architecture' in title5.lower()

        autoshapes5 = get_autoshapes_with_text(slide5)
        texts5 = {t.lower() for t, _ in autoshapes5 if t}
        data_boxes = {'sql db', 'nosql db', 'file storage'}
        has_client5 = client_boxes.issubset(texts5)
        has_server5 = server_boxes.issubset(texts5)
        has_data5 = data_boxes.issubset(texts5)

        if has_title5 and has_client5 and has_server5 and has_data5:
            print(f"PASS: Component 3 - Slide 5 has title '{title5}' and all 3 layers (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 - Slide 5: title='{title5}' (ok={has_title5}), "
                  f"client={has_client5}, server={has_server5}, data={has_data5}, texts={texts5}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # --- Component 4: Color differentiation per layer (0.15 points) ---
    # Blue for client, green for server, orange for data on slide 5 (has all layers)
    try:
        color_score = 0.0
        client_colors = []
        server_colors = []
        data_colors = []

        for text, fill_color in autoshapes5:
            text_lower = text.lower()
            if text_lower in client_boxes:
                client_colors.append(classify_color(fill_color))
            elif text_lower in server_boxes:
                server_colors.append(classify_color(fill_color))
            elif text_lower in data_boxes:
                data_colors.append(classify_color(fill_color))

        # Each layer gets 0.05 for correct color
        if client_colors and all(c == 'blue' for c in client_colors):
            color_score += 0.05
            print(f"  Client layer: blue colors confirmed")
        else:
            print(f"  Client layer: expected blue, got {client_colors}")

        if server_colors and all(c == 'green' for c in server_colors):
            color_score += 0.05
            print(f"  Server layer: green colors confirmed")
        else:
            print(f"  Server layer: expected green, got {server_colors}")

        if data_colors and all(c == 'orange' for c in data_colors):
            color_score += 0.05
            print(f"  Data layer: orange colors confirmed")
        else:
            print(f"  Data layer: expected orange, got {data_colors}")

        if color_score > 0:
            print(f"PASS: Component 4 - Color differentiation ({color_score} pts)")
            total_score += color_score
        else:
            print(f"FAIL: Component 4 - No correct layer colors found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # --- Component 5: Connecting arrows between layers (0.10 points) ---
    # Arrows are autoshapes without text that connect layers
    # Slide 4 should have arrows between client and server (at least 1)
    # Slide 5 should have arrows between all layers (more than slide 4)
    try:
        empty_shapes4 = [s for s in autoshapes4 if not s[0]]  # no text = arrow/connector
        empty_shapes5 = [s for s in autoshapes5 if not s[0]]

        arrow_score = 0.0
        if len(empty_shapes4) >= 1:
            arrow_score += 0.05
            print(f"  Slide 4: {len(empty_shapes4)} connecting shapes found")
        else:
            print(f"  Slide 4: no connecting shapes found")

        if len(empty_shapes5) >= 2 and len(empty_shapes5) > len(empty_shapes4):
            arrow_score += 0.05
            print(f"  Slide 5: {len(empty_shapes5)} connecting shapes found (more than slide 4)")
        elif len(empty_shapes5) >= 2:
            arrow_score += 0.03
            print(f"  Slide 5: {len(empty_shapes5)} connecting shapes found (same or fewer than slide 4)")
        else:
            print(f"  Slide 5: insufficient connecting shapes ({len(empty_shapes5)})")

        if arrow_score > 0:
            print(f"PASS: Component 5 - Arrows ({arrow_score} pts)")
            total_score += arrow_score
        else:
            print(f"FAIL: Component 5 - No connecting arrows found")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # --- Component 6: Morph transitions on slides 4 and 5 (0.15 points) ---
    try:
        morph4 = check_morph_transition(file_path, 4)
        morph5 = check_morph_transition(file_path, 5)

        morph_score = 0.0
        if morph4:
            morph_score += 0.075
            print(f"  Slide 4: Morph transition found")
        else:
            print(f"  Slide 4: No Morph transition")

        if morph5:
            morph_score += 0.075
            print(f"  Slide 5: Morph transition found")
        else:
            print(f"  Slide 5: No Morph transition")

        if morph_score > 0:
            print(f"PASS: Component 6 - Morph transitions ({morph_score} pts)")
            total_score += morph_score
        else:
            print(f"FAIL: Component 6 - No Morph transitions found")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
