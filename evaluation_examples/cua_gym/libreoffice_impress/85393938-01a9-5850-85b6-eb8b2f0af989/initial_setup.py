"""
Initial Setup: Multi-slide build animation for system architecture
Task ID: impress_stu_090
Domain: libreoffice_impress

Creates an 8-slide presentation 'System_Design.pptx' where slides 3-5 are
intentionally left empty (blank) for the agent to populate with architecture
layer diagrams, connecting arrows, colored boxes, and Morph transitions.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_090'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text):
    """Add a title textbox at the top of a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x2E, 0x2E)


def add_body_text(slide, body_text, top=Inches(1.2)):
    """Add body text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.7), top, Inches(8.5), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only
    add_title_text(slide1, "System Architecture Overview")
    txBox = slide1.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.5), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cloud-Native Microservices Platform"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2 = tf.add_paragraph()
    p2.text = "Engineering Division | Q2 2025 Technical Review"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    # --- Slide 2: Project Background ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_text(slide2, "Project Background")
    add_body_text(slide2, [
        "Our platform serves 2.4 million daily active users across 14 regions.",
        "Current monolithic architecture creates deployment bottlenecks and scaling issues.",
        "Migration to microservices began in January 2025 with a phased rollout plan.",
        "Key drivers: horizontal scalability, independent deployments, fault isolation.",
        "Target completion: December 2025 with full production migration.",
    ])

    # --- Slides 3-5: EMPTY (blank) ---
    # These slides must remain empty for the agent to populate
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])  # Title Only (blank-like)

    # --- Slide 6: Timeline ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_text(slide6, "Migration Timeline")
    add_body_text(slide6, [
        "Phase 1 (Jan-Mar): Service decomposition and API gateway setup",
        "Phase 2 (Apr-Jun): Authentication and core logic migration",
        "Phase 3 (Jul-Sep): Data layer migration and caching strategy",
        "Phase 4 (Oct-Dec): Performance optimization and full cutover",
        "Rollback windows: 48 hours per phase, full rollback capability until Phase 4",
    ])

    # --- Slide 7: Team Responsibilities ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_text(slide7, "Team Responsibilities")
    add_body_text(slide7, [
        "Platform Team (Elena Rodriguez): Infrastructure, CI/CD, monitoring",
        "Backend Team (James Park): Service decomposition, API contracts",
        "Data Team (Priya Sharma): Database migration, caching, replication",
        "Frontend Team (Alex Nguyen): Client SDK updates, progressive rollout",
        "SRE Team (David Kim): Reliability targets, incident response, load testing",
    ])

    # --- Slide 8: Q&A ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_text(slide8, "Questions & Discussion")
    txBox = slide8.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank you for your attention."
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2 = tf.add_paragraph()
    p2.text = "Contact: architecture-team@company.com"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
