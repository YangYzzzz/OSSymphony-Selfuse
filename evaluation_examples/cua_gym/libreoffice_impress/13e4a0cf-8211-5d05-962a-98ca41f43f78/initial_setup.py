"""
Initial Setup: Create a 4-slide Lecture presentation with no speaker notes.
Task ID: impress_wf_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_009'
OUTPUT = f'{WORKDIR}/Desktop/Lecture.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Introduction to Machine Learning"
    slide1.placeholders[1].text = "CS 4501 — Spring 2025\nDr. Elena Rodriguez"
    # Set title font
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # --- Slide 2: Key Concepts ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Concepts in Supervised Learning"
    body2 = slide2.placeholders[1].text_frame
    body2.clear()
    concepts = [
        "Classification vs. Regression: predicting categories versus continuous values",
        "Training Set: labeled examples used to learn model parameters",
        "Validation Set: held-out data for tuning hyperparameters",
        "Test Set: unseen data for final performance evaluation",
        "Overfitting: when a model memorizes training data but fails to generalize",
        "Bias-Variance Tradeoff: balancing underfitting and overfitting",
    ]
    for i, concept in enumerate(concepts):
        if i == 0:
            body2.paragraphs[0].text = concept
        else:
            p = body2.add_paragraph()
            p.text = concept
            p.level = 0

    # --- Slide 3: Live Demo ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Live Demo: Training a Decision Tree"
    body3 = slide3.placeholders[1].text_frame
    body3.clear()
    demo_items = [
        "Dataset: Iris flower classification (150 samples, 4 features)",
        "Tool: scikit-learn in a Jupyter Notebook",
        "Steps: load data → split → fit model → visualize tree → evaluate accuracy",
        "Expected accuracy: ~95% on test set",
        "Discussion: how does tree depth affect performance?",
    ]
    for i, item in enumerate(demo_items):
        if i == 0:
            body3.paragraphs[0].text = item
        else:
            p = body3.add_paragraph()
            p.text = item
            p.level = 0

    # --- Slide 4: Summary & Q&A ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Summary & Questions"
    body4 = slide4.placeholders[1].text_frame
    body4.clear()
    summary_items = [
        "Supervised learning requires labeled training data",
        "Model selection depends on task type and data characteristics",
        "Always evaluate on held-out test data to measure generalization",
        "Next lecture: Unsupervised Learning and Clustering",
        "",
        "Questions? Office hours: Tue/Thu 3-4 PM, Room 312",
    ]
    for i, item in enumerate(summary_items):
        if i == 0:
            body4.paragraphs[0].text = item
        else:
            p = body4.add_paragraph()
            p.text = item
            p.level = 0

    # IMPORTANT: Do NOT access notes_slide — that would create notes
    # The task requires no speaker notes initially

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
