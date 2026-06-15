"""
Initial Setup: Product Backlog Presentation — pre-task state
Task ID: osworld_impress_strikethrough_text_012
Domain: libreoffice_impress

Creates a 6-slide product backlog deck.
Slide 4 has a content textbox with 5 feature items in plain black text.
NO strikethrough and NO gray color on any item (that is the task to perform).
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_slide_title_content(prs, title_text, content_items, bullet_color=None):
    """Add a slide using Title+Content layout (index 1) with given title and bullet items."""
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_run = title_shape.text_frame.paragraphs[0].runs[0]
    title_run.font.bold = True
    title_run.font.size = Pt(32)
    title_run.font.color.rgb = RGBColor(0x1F, 0x38, 0x96)  # dark blue

    # Set content placeholder
    content_ph = slide.placeholders[1]
    tf = content_ph.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = item
        para.level = 0
        run = para.runs[0]
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # black
        run.font.bold = False

    return slide


def add_title_slide(prs, title_text, subtitle_text):
    """Add a title slide (layout 0)."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text

    title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.bold = True
    title_run.font.size = Pt(36)
    title_run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    sub_run = slide.placeholders[1].text_frame.paragraphs[0].runs[0]
    sub_run.font.size = Pt(22)
    sub_run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(
        prs,
        "Product Backlog Q2 2025",
        "Engineering Team — Sprint Planning"
    )

    # Slide 2: Overview
    add_slide_title_content(prs, "Backlog Overview", [
        "Total items: 47 across 5 epics",
        "Sprint velocity: 32 story points",
        "Carry-over from Q1: 8 items",
        "New requests from stakeholders: 12 items",
        "Blocked items: 3 (dependency on API team)",
    ])

    # Slide 3: High Priority Items
    add_slide_title_content(prs, "High Priority — Must Have", [
        "User authentication with SSO (8 pts)",
        "Dashboard performance optimization (13 pts)",
        "CSV export for all report types (5 pts)",
        "Mobile-responsive navigation redesign (8 pts)",
        "GDPR data deletion workflow (5 pts)",
        "API rate limiting and throttling (3 pts)",
    ])

    # Slide 4: Feature Backlog — THE KEY SLIDE
    # 5 items in black text, NO strikethrough, NO gray — this is what the agent must modify
    feature_items = [
        "Implement real-time notifications via WebSocket",
        "Refactor legacy payment module to Stripe v3",
        "Add dark mode support across all UI components",
        "Integrate Salesforce CRM data sync pipeline",
        "Improve search indexing with Elasticsearch upgrade",
    ]
    add_slide_title_content(prs, "Feature Backlog — In Progress", feature_items)

    # Slide 5: Technical Debt
    add_slide_title_content(prs, "Technical Debt Items", [
        "Migrate database from PostgreSQL 12 to 15",
        "Remove deprecated jQuery 1.x dependencies",
        "Consolidate microservices config management",
        "Update Docker base images to Ubuntu 22.04 LTS",
        "Add integration test coverage for billing module",
    ])

    # Slide 6: Done / Released
    add_slide_title_content(prs, "Recently Completed", [
        "Two-factor authentication via SMS and TOTP",
        "Webhook delivery retry mechanism with backoff",
        "Automated nightly backups to S3 with encryption",
        "Onboarding wizard for new enterprise customers",
        "Accessibility audit fixes (WCAG 2.1 AA compliance)",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
