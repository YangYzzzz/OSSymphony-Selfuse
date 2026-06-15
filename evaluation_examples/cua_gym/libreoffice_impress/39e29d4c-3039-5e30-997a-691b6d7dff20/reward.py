"""
Reward Script: Apply strikethrough to items 1 and 2 on slide 4, change their color to gray,
               and leave items 3, 4, and 5 with no strikethrough and in their original black color.
Task ID: osworld_impress_strikethrough_text_012
Domain: libreoffice_impress
Scoring:
  Component 1: Item 1 (para 0) has strikethrough (sngStrike)    — 0.25 pts
  Component 2: Item 1 (para 0) has gray color (808080)           — 0.25 pts
  Component 3: Item 2 (para 1) has strikethrough (sngStrike)    — 0.25 pts
  Component 4: Item 2 (para 1) has gray color (808080)           — 0.25 pts
  Total: 1.0

  Items 3-5 (para 2-4) should be noStrike and black (000000) — verified as preconditions/gates.
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_012'

# Gray color expected for items 1 and 2 after task completion
GRAY_COLOR = '808080'
# Black color expected for items 3-5 (unchanged)
BLACK_COLOR = '000000'


def get_run_strike(run):
    """Return the strike attribute from the run XML element."""
    return run.font._element.attrib.get('strike', 'noStrike')


def get_run_rgb(run):
    """Return the RGB hex string for a run's font color, or None if not set."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
        return None
    except Exception:
        return None


def get_content_placeholder(slide):
    """Return the 'Content Placeholder 2' shape on the given slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == 'Content Placeholder 2':
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation — gate
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify correct number of slides — gate
    if len(prs.slides) != 6:
        print(f"CRITICAL: Expected 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide 4 (0-indexed: index 3)
    slide4 = prs.slides[3]

    # Find the content placeholder on slide 4
    content_shape = get_content_placeholder(slide4)
    if content_shape is None:
        print("CRITICAL: 'Content Placeholder 2' not found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = content_shape.text_frame.paragraphs

    # Verify there are at least 5 items — gate
    if len(paragraphs) < 5:
        print(f"CRITICAL: Expected at least 5 paragraphs in content placeholder, found {len(paragraphs)}")
        print("REWARD: 0.0")
        return 0.0

    # Helper: get first non-empty run in a paragraph
    def get_first_run(para):
        runs = [r for r in para.runs if (r.text or '').strip()]
        return runs[0] if runs else (para.runs[0] if para.runs else None)

    # Component 1: Item 1 (para index 0) has strikethrough sngStrike (0.25 points)
    try:
        run1 = get_first_run(paragraphs[0])
        if run1 is None:
            print("FAIL: Component 1 — Para 0 has no runs")
        else:
            strike1 = get_run_strike(run1)
            if strike1 == 'sngStrike':
                print(f"PASS: Component 1 — Item 1 has sngStrike (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Item 1 strike expected 'sngStrike', found '{strike1}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Item 1 (para index 0) has gray color 808080 (0.25 points)
    try:
        run1 = get_first_run(paragraphs[0])
        if run1 is None:
            print("FAIL: Component 2 — Para 0 has no runs")
        else:
            rgb1 = get_run_rgb(run1)
            if rgb1 is not None and rgb1.upper() == GRAY_COLOR.upper():
                print(f"PASS: Component 2 — Item 1 color is gray 808080 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Item 1 color expected '{GRAY_COLOR}', found '{rgb1}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Item 2 (para index 1) has strikethrough sngStrike (0.25 points)
    try:
        run2 = get_first_run(paragraphs[1])
        if run2 is None:
            print("FAIL: Component 3 — Para 1 has no runs")
        else:
            strike2 = get_run_strike(run2)
            if strike2 == 'sngStrike':
                print(f"PASS: Component 3 — Item 2 has sngStrike (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 — Item 2 strike expected 'sngStrike', found '{strike2}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Item 2 (para index 1) has gray color 808080 (0.25 points)
    try:
        run2 = get_first_run(paragraphs[1])
        if run2 is None:
            print("FAIL: Component 4 — Para 1 has no runs")
        else:
            rgb2 = get_run_rgb(run2)
            if rgb2 is not None and rgb2.upper() == GRAY_COLOR.upper():
                print(f"PASS: Component 4 — Item 2 color is gray 808080 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 — Item 2 color expected '{GRAY_COLOR}', found '{rgb2}'")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Precondition gate: items 3-5 should remain unchanged (noStrike, black)
    # Not scored, but logged for debugging
    unchanged_warn_count = 0
    for idx in [2, 3, 4]:
        try:
            run = get_first_run(paragraphs[idx])
            if run is None:
                continue
            strike = get_run_strike(run)
            rgb = get_run_rgb(run)
            if strike != 'noStrike':
                print(f"WARN: Item {idx+1} (para {idx}) has unexpected strike '{strike}' — should be noStrike")
                unchanged_warn_count += 1
            if rgb is not None and rgb.upper() != BLACK_COLOR.upper():
                print(f"WARN: Item {idx+1} (para {idx}) has unexpected color '{rgb}' — should be black or None")
        except Exception as e:
            print(f"WARN: Could not verify item {idx+1}: {e}")

    if unchanged_warn_count == 0:
        print("INFO: Items 3-5 remain unchanged (noStrike, black) — precondition satisfied")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
