"""
Initial Setup: Create a 12-slide ENV 200 Environmental Policy presentation with no footer elements.
Task ID: impress_stu_091
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_091'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Color palette
    green_dark = RGBColor(0x1B, 0x5E, 0x20)
    green_medium = RGBColor(0x2E, 0x7D, 0x32)
    green_light = RGBColor(0x4C, 0xAF, 0x50)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    dark_text = RGBColor(0x33, 0x33, 0x33)
    blue_accent = RGBColor(0x00, 0x27, 0x52)

    # ========== SLIDE 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = green_dark

    add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                "ENV 200 - Environmental Policy", font_size=40, bold=True,
                alignment=PP_ALIGN.CENTER, color=white)
    add_textbox(slide1, Inches(1), Inches(3.2), Inches(11), Inches(1),
                "An Introduction to Environmental Law, Regulation, and Governance",
                font_size=22, alignment=PP_ALIGN.CENTER, color=white)
    add_textbox(slide1, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
                "University of California, Davis  |  Spring 2025",
                font_size=16, alignment=PP_ALIGN.CENTER, color=white)
    add_textbox(slide1, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                "Professor Rachel M. Torres, Department of Environmental Science & Policy",
                font_size=14, alignment=PP_ALIGN.CENTER, color=white)

    # ========== SLIDE 2: Course Overview ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Course Overview", font_size=32, bold=True, color=green_dark)
    add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "This course examines the development and implementation of environmental policy "
                "in the United States and internationally. Students will explore the intersection "
                "of science, law, economics, and politics in shaping environmental governance.\n\n"
                "Key themes include:\n"
                "• The evolution of environmental legislation from the 1970s to present\n"
                "• Cost-benefit analysis in environmental regulation\n"
                "• The role of federal agencies (EPA, NOAA, DOI) in policy enforcement\n"
                "• International environmental agreements and their effectiveness\n"
                "• Environmental justice and equity considerations",
                font_size=16, color=dark_text)

    # ========== SLIDE 3: History of Environmental Law ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "History of Environmental Legislation", font_size=32, bold=True, color=green_dark)
    add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
                "Landmark Environmental Laws:\n\n"
                "1969 - National Environmental Policy Act (NEPA)\n"
                "1970 - Clean Air Act (CAA)\n"
                "1972 - Clean Water Act (CWA)\n"
                "1973 - Endangered Species Act (ESA)\n"
                "1976 - Resource Conservation and Recovery Act (RCRA)\n"
                "1980 - Comprehensive Environmental Response Act (CERCLA)\n"
                "1990 - Clean Air Act Amendments\n"
                "2015 - Paris Climate Agreement (International)",
                font_size=14, color=dark_text)
    add_textbox(slide3, Inches(7), Inches(1.8), Inches(5.5), Inches(5),
                "The 'Environmental Decade' of the 1970s saw unprecedented legislative action. "
                "President Nixon established the EPA in 1970, consolidating federal environmental "
                "responsibilities under one agency. These foundational laws established the framework "
                "for environmental protection that persists today, though implementation and enforcement "
                "have varied significantly across administrations.",
                font_size=14, color=dark_text)

    # ========== SLIDE 4: Clean Air Act ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "The Clean Air Act: A Case Study", font_size=32, bold=True, color=green_dark)
    add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "The Clean Air Act (1970, amended 1977, 1990) is one of the most comprehensive "
                "air quality laws in the world.\n\n"
                "Key Provisions:\n"
                "• National Ambient Air Quality Standards (NAAQS) for six criteria pollutants\n"
                "• New Source Performance Standards (NSPS) for industrial facilities\n"
                "• Hazardous Air Pollutant regulations (Section 112)\n"
                "• Title IV Acid Rain Program - first major cap-and-trade system\n"
                "• Mobile source emission standards for vehicles\n\n"
                "Economic Impact: EPA estimates the 1990 amendments produced benefits of $2 trillion "
                "against costs of $65 billion by 2020 - a benefit-cost ratio of approximately 30:1.",
                font_size=14, color=dark_text)

    # ========== SLIDE 5: Regulatory Framework ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "The Regulatory Framework", font_size=32, bold=True, color=green_dark)
    add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                "Command-and-Control Regulation:\n\n"
                "• Technology-based standards\n"
                "• Performance-based standards\n"
                "• Ambient quality standards\n"
                "• Permits and licensing requirements\n"
                "• Monitoring and reporting obligations\n\n"
                "Advantages: Certainty of outcome, clear compliance criteria\n"
                "Disadvantages: Inflexible, potentially high compliance costs",
                font_size=14, color=dark_text)
    add_textbox(slide5, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
                "Market-Based Instruments:\n\n"
                "• Cap-and-trade systems (SO₂, CO₂)\n"
                "• Environmental taxes and fees\n"
                "• Tradeable permit programs\n"
                "• Deposit-refund systems\n"
                "• Subsidies for green technology\n\n"
                "Advantages: Cost-effective, encourages innovation\n"
                "Disadvantages: Uncertain environmental outcomes, political challenges",
                font_size=14, color=dark_text)

    # ========== SLIDE 6: Cost-Benefit Analysis ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Cost-Benefit Analysis in Environmental Policy", font_size=32, bold=True, color=green_dark)
    add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "Since Executive Order 12866 (1993), major federal regulations require a regulatory "
                "impact analysis demonstrating that benefits justify costs.\n\n"
                "Challenges in Environmental CBA:\n"
                "• Valuing non-market goods (clean air, biodiversity, ecosystem services)\n"
                "• Discount rates for long-term impacts (climate change damages)\n"
                "• Distributional effects across income groups and communities\n"
                "• Scientific uncertainty in dose-response relationships\n"
                "• Co-benefits vs. targeted benefits of regulations\n\n"
                "The Social Cost of Carbon (SCC) has been particularly contentious, ranging from "
                "$1/ton to over $200/ton depending on discount rate and modeling assumptions. "
                "The Biden administration set the interim SCC at $51/ton in 2021.",
                font_size=14, color=dark_text)

    # ========== SLIDE 7: Environmental Justice ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Environmental Justice", font_size=32, bold=True, color=green_dark)
    add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "Environmental justice addresses the disproportionate environmental burdens borne by "
                "low-income communities and communities of color.\n\n"
                "Key Milestones:\n"
                "• 1982 - Warren County PCB landfill protests (North Carolina)\n"
                "• 1987 - UCC report 'Toxic Wastes and Race in the United States'\n"
                "• 1994 - Executive Order 12898 (Federal Actions to Address EJ)\n"
                "• 2021 - Justice40 Initiative (40% of federal climate investments to disadvantaged communities)\n\n"
                "Research consistently shows that hazardous waste facilities, polluting industries, "
                "and transportation corridors are disproportionately located near minority and "
                "low-income neighborhoods. The Flint water crisis (2014-2019) remains a prominent "
                "example of environmental injustice in infrastructure management.",
                font_size=14, color=dark_text)

    # ========== SLIDE 8: Climate Policy ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "U.S. Climate Change Policy", font_size=32, bold=True, color=green_dark)
    add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "The U.S. approach to climate change has oscillated significantly between administrations.\n\n"
                "Timeline of Key Actions:\n"
                "• 1992 - UNFCCC ratification (non-binding emissions framework)\n"
                "• 1997 - Kyoto Protocol signed but never ratified by U.S. Senate\n"
                "• 2007 - Massachusetts v. EPA (Supreme Court: CO₂ is a pollutant under CAA)\n"
                "• 2009 - EPA Endangerment Finding for greenhouse gases\n"
                "• 2015 - Clean Power Plan proposed (stayed by Supreme Court 2016)\n"
                "• 2015 - Paris Agreement signed (U.S. withdrew 2020, rejoined 2021)\n"
                "• 2022 - Inflation Reduction Act ($369B in climate and energy provisions)\n\n"
                "Current U.S. NDC target: 50-52% reduction in GHG emissions below 2005 levels by 2030.",
                font_size=14, color=dark_text)

    # ========== SLIDE 9: International Agreements ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "International Environmental Agreements", font_size=32, bold=True, color=green_dark)
    add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "Major Multilateral Environmental Agreements:\n\n"
                "• Montreal Protocol (1987) - Ozone layer protection; widely considered the most "
                "successful international environmental treaty. Phased out 99% of ODS production.\n"
                "• Basel Convention (1989) - Transboundary movement of hazardous wastes\n"
                "• Convention on Biological Diversity (1992) - Biodiversity conservation\n"
                "• Kyoto Protocol (1997) - Binding GHG reduction targets for developed nations\n"
                "• Stockholm Convention (2001) - Persistent organic pollutants (POPs)\n"
                "• Paris Agreement (2015) - Universal climate framework with NDCs\n"
                "• Kunming-Montreal Framework (2022) - '30x30' biodiversity protection target\n\n"
                "The success of the Montreal Protocol offers lessons: clear science, available "
                "substitutes, phased implementation, and financial mechanisms for developing countries.",
                font_size=14, color=dark_text)

    # ========== SLIDE 10: Water Policy ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide10, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Water Resource Management", font_size=32, bold=True, color=green_dark)
    add_textbox(slide10, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "Water policy operates at the intersection of federal, state, and local authority.\n\n"
                "Federal Framework:\n"
                "• Clean Water Act - Point source discharge permits (NPDES), water quality standards\n"
                "• Safe Drinking Water Act - Public water system standards, underground injection control\n"
                "• Waters of the United States (WOTUS) - Ongoing jurisdictional debate\n\n"
                "California Water Issues (Relevant to UC Davis):\n"
                "• Central Valley agriculture consumes ~80% of developed water supply\n"
                "• Sacramento-San Joaquin Delta ecosystem vs. water supply tensions\n"
                "• Sustainable Groundwater Management Act (SGMA, 2014)\n"
                "• Drought emergency declarations and mandatory conservation measures\n"
                "• Colorado River compact renegotiation and 'Tier 1' shortage conditions",
                font_size=14, color=dark_text)

    # ========== SLIDE 11: Emerging Issues ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide11, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Emerging Environmental Policy Challenges", font_size=32, bold=True, color=green_dark)
    add_textbox(slide11, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                "New Frontiers in Environmental Regulation:\n\n"
                "• PFAS ('forever chemicals') - EPA proposed NPDWR for 6 PFAS compounds in 2023\n"
                "• Microplastics - No federal regulatory framework yet; California leads with SB 1422\n"
                "• Environmental impacts of AI and data centers (energy, water consumption)\n"
                "• Permitting reform for renewable energy infrastructure\n"
                "• Carbon capture and storage (CCS) liability and regulatory frameworks\n"
                "• Climate litigation - Juliana v. United States, state-level lawsuits against fossil fuel companies\n"
                "• Critical minerals for clean energy transition - environmental impacts of lithium, cobalt mining\n"
                "• Agricultural runoff and nonpoint source pollution (largely unregulated under CWA)",
                font_size=14, color=dark_text)

    # ========== SLIDE 12: Course Assignments ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide12, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                "Course Assignments & Grading", font_size=32, bold=True, color=green_dark)
    add_textbox(slide12, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                "Assignment Breakdown:\n\n"
                "• Weekly Discussion Posts: 15%\n"
                "• Policy Brief (4-5 pages): 20%\n"
                "• Regulatory Analysis Project: 25%\n"
                "• Midterm Examination: 15%\n"
                "• Final Presentation: 15%\n"
                "• Participation: 10%\n\n"
                "Total: 100%",
                font_size=14, color=dark_text)
    add_textbox(slide12, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5),
                "Important Dates:\n\n"
                "• Week 3: Discussion groups assigned\n"
                "• Week 5: Policy brief topic proposal due\n"
                "• Week 7: Midterm examination\n"
                "• Week 8: Policy brief due\n"
                "• Week 10: Regulatory analysis project due\n"
                "• Finals Week: Group presentations\n\n"
                "Office Hours: Tues/Thurs 2:00-3:30 PM, 2138 Wickson Hall",
                font_size=14, color=dark_text)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
