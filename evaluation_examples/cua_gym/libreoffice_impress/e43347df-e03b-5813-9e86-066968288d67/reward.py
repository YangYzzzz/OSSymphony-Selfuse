"""
Reward Script: Custom footer bar on slides 2-12
Task ID: impress_stu_091
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Dark charcoal rectangle bars on slides 2-12
  Component 2 (0.25): Left-aligned course name text in footer
  Component 3 (0.25): Center-aligned university name text in footer
  Component 4 (0.15): Right-aligned slide numbers in footer
  Component 5 (0.10): Slide 1 has no footer bar
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_091'

# Constants from task spec
EXPECTED_FILL_RGB = '2D2D2D'   # dark charcoal
EXPECTED_FONT_COLOR = 'FFFFFF'  # white
EXPECTED_FONT_SIZE = Pt(9)      # 114300 EMU
BAR_HEIGHT_EMU = int(0.4 * 914400)  # 365760 EMU
FOOTER_SLIDES = list(range(1, 12))  # 0-indexed slides 1..11 = presentation slides 2..12


def has_footer_bar(slide, slide_width):
    """Check if slide has a rectangle shape at the bottom with charcoal fill, full width, 0.4in tall."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check position: at or near bottom, full width, correct height
            width_ok = abs(shape.width - slide_width) / max(shape.width, slide_width) <= 0.01
            height_ok = abs(shape.height - BAR_HEIGHT_EMU) / max(shape.height, BAR_HEIGHT_EMU) <= 0.05
            # Check fill color
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    color_ok = str(fill.fore_color.rgb).upper() == EXPECTED_FILL_RGB.upper()
                else:
                    color_ok = False
            except Exception:
                color_ok = False

            if width_ok and height_ok and color_ok:
                return True
    return False


def find_footer_texts(slide):
    """Find text boxes in the footer region (bottom area) of a slide.
    Returns dict with keys 'left', 'center', 'right' based on alignment."""
    slide_height = 6858000  # 7.5 inches
    footer_region_top = slide_height - int(0.6 * 914400)  # bottom 0.6 inches

    result = {'left': None, 'center': None, 'right': None}

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
            # Only consider shapes in the footer region
            if shape.top >= footer_region_top:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    alignment = para.alignment
                    # Normalize None to LEFT
                    if alignment is None or alignment == PP_ALIGN.LEFT:
                        key = 'left'
                    elif alignment == PP_ALIGN.CENTER:
                        key = 'center'
                    elif alignment == PP_ALIGN.RIGHT:
                        key = 'right'
                    else:
                        key = 'left'

                    # Get font properties from first run
                    font_color = None
                    font_size = None
                    if para.runs:
                        run = para.runs[0]
                        try:
                            font_color = str(run.font.color.rgb).upper()
                        except Exception:
                            font_color = None
                        font_size = run.font.size

                    result[key] = {
                        'text': text,
                        'font_color': font_color,
                        'font_size': font_size,
                    }
    return result


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 12:
        print(f"CRITICAL: Expected at least 12 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width

    # Component 1: Dark charcoal rectangle bars on slides 2-12 (0.25 points)
    # Each of the 11 slides contributes equally
    try:
        bars_found = 0
        for idx in FOOTER_SLIDES:
            slide = prs.slides[idx]
            if has_footer_bar(slide, slide_width):
                bars_found += 1

        if bars_found == len(FOOTER_SLIDES):
            print(f"PASS: Component 1 - All {bars_found}/11 slides have charcoal footer bar (0.25 pts)")
            total_score += 0.25
        elif bars_found > 0:
            partial = 0.25 * (bars_found / len(FOOTER_SLIDES))
            print(f"PARTIAL: Component 1 - {bars_found}/11 slides have charcoal footer bar ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No slides have charcoal footer bar")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Left-aligned 'ENV 200 - Environmental Policy' in 9pt white (0.25 points)
    try:
        left_ok = 0
        for idx in FOOTER_SLIDES:
            slide = prs.slides[idx]
            footer = find_footer_texts(slide)
            left_info = footer.get('left')
            if left_info is not None:
                text_match = 'ENV 200' in left_info['text'] and 'Environmental Policy' in left_info['text']
                color_match = left_info['font_color'] == EXPECTED_FONT_COLOR
                size_match = left_info['font_size'] is not None and abs(left_info['font_size'] - EXPECTED_FONT_SIZE) <= 12700  # 1pt tolerance
                if text_match and color_match and size_match:
                    left_ok += 1

        if left_ok == len(FOOTER_SLIDES):
            print(f"PASS: Component 2 - All {left_ok}/11 slides have correct left footer text (0.25 pts)")
            total_score += 0.25
        elif left_ok > 0:
            partial = 0.25 * (left_ok / len(FOOTER_SLIDES))
            print(f"PARTIAL: Component 2 - {left_ok}/11 slides have correct left footer text ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No slides have correct left footer text")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Center-aligned 'University of California, Davis' in 9pt white (0.25 points)
    try:
        center_ok = 0
        for idx in FOOTER_SLIDES:
            slide = prs.slides[idx]
            footer = find_footer_texts(slide)
            center_info = footer.get('center')
            if center_info is not None:
                text_match = 'University of California' in center_info['text'] and 'Davis' in center_info['text']
                color_match = center_info['font_color'] == EXPECTED_FONT_COLOR
                size_match = center_info['font_size'] is not None and abs(center_info['font_size'] - EXPECTED_FONT_SIZE) <= 12700
                if text_match and color_match and size_match:
                    center_ok += 1

        if center_ok == len(FOOTER_SLIDES):
            print(f"PASS: Component 3 - All {center_ok}/11 slides have correct center footer text (0.25 pts)")
            total_score += 0.25
        elif center_ok > 0:
            partial = 0.25 * (center_ok / len(FOOTER_SLIDES))
            print(f"PARTIAL: Component 3 - {center_ok}/11 slides have correct center footer text ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No slides have correct center footer text")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Right-aligned slide number in 9pt white (0.15 points)
    try:
        right_ok = 0
        for idx in FOOTER_SLIDES:
            slide = prs.slides[idx]
            footer = find_footer_texts(slide)
            right_info = footer.get('right')
            if right_info is not None:
                # Slide number should match (idx+1 is the 1-based slide number)
                expected_num = str(idx + 1)
                text_match = right_info['text'].strip() == expected_num
                color_match = right_info['font_color'] == EXPECTED_FONT_COLOR
                size_match = right_info['font_size'] is not None and abs(right_info['font_size'] - EXPECTED_FONT_SIZE) <= 12700
                if text_match and color_match and size_match:
                    right_ok += 1

        if right_ok == len(FOOTER_SLIDES):
            print(f"PASS: Component 4 - All {right_ok}/11 slides have correct slide number (0.15 pts)")
            total_score += 0.15
        elif right_ok > 0:
            partial = 0.15 * (right_ok / len(FOOTER_SLIDES))
            print(f"PARTIAL: Component 4 - {right_ok}/11 slides have correct slide number ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No slides have correct slide number")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Slide 1 has NO footer bar while other slides DO (0.10 points)
    # Compound check: slide 1 must lack footer AND at least one other slide must have one.
    # This ensures we only score the intentional exclusion of slide 1, not the pre-task state.
    try:
        slide1 = prs.slides[0]
        slide1_no_bar = not has_footer_bar(slide1, slide_width)
        any_other_has_bar = bars_found > 0  # from Component 1
        if slide1_no_bar and any_other_has_bar:
            print(f"PASS: Component 5 - Slide 1 correctly excluded from footer bars (0.10 pts)")
            total_score += 0.10
        elif not slide1_no_bar:
            print(f"FAIL: Component 5 - Slide 1 has a footer bar (should not)")
        else:
            print(f"FAIL: Component 5 - No other slides have footer bars yet, cannot verify exclusion")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
