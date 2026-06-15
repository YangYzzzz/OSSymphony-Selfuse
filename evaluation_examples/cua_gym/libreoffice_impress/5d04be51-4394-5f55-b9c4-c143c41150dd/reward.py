"""
Reward Script: Neuroscience Research Poster (48x36 inches)
Task ID: impress_stu_062
Domain: libreoffice_impress
Scoring:
  C1: Header bar with #003366 fill (0.15)
  C2: Title text correct, 54pt bold white (0.15)
  C3: Author line present, 24pt white (0.10)
  C4: Three-column layout with 7 required sections (0.20)
  C5: Section headers 28pt bold #003366 (0.15)
  C6: Col 1 backgrounds #E8F0FE (0.10)
  C7: Placeholders for brain scan image and bar chart (0.10)
  C8: University logo placeholder in top-right area (0.05)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_062'


def get_all_shapes(slide):
    """Recursively extract all shapes, including grouped ones."""
    results = []
    for shape in slide.shapes:
        results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.append(sub)
    return results


def get_shape_fill_rgb(shape):
    """Get solid fill RGB color from a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_run_color_rgb(run):
    """Get run font color as hex string, or None."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def find_text_shapes_containing(shapes, text_lower):
    """Find shapes whose text contains the given substring (case-insensitive)."""
    results = []
    for shape in shapes:
        if hasattr(shape, 'text') and text_lower in (shape.text or '').lower():
            results.append(shape)
    return results


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: single slide, 48x36 inches
    if len(prs.slides) < 1:
        print("FAIL: No slides found")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]
    all_shapes = get_all_shapes(slide)

    # Check slide dimensions are approximately 48x36
    width_in = prs.slide_width / 914400
    height_in = prs.slide_height / 914400
    if abs(width_in - 48) > 1 or abs(height_in - 36) > 1:
        print(f"FAIL: Slide dimensions {width_in}x{height_in} not approximately 48x36 inches")
        print("REWARD: 0.0")
        return 0.0

    # Must have meaningful content (more than just a blank slide)
    if len(all_shapes) < 5:
        print(f"FAIL: Only {len(all_shapes)} shapes found, expected poster layout with many shapes")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Header bar with dark blue (#003366) fill (0.15 points)
    try:
        header_bar_found = False
        for shape in all_shapes:
            # Look for a rectangle at the top of the slide with #003366 fill
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                fill_rgb = get_shape_fill_rgb(shape)
                if fill_rgb and fill_rgb.upper() == '003366':
                    # Should span most of the width and be at the top
                    if shape.top < Inches(2) and shape.width > Inches(30):
                        header_bar_found = True
                        break
        if header_bar_found:
            print(f"PASS: Component 1 -- Header bar with #003366 fill found at top (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- No dark blue (#003366) header bar found at top of slide")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Title text 'Neural Correlates of Decision-Making Under Uncertainty', 54pt bold white (0.15 points)
    try:
        title_found = False
        for shape in all_shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            text = (shape.text or '').strip()
            if 'neural correlates' in text.lower() and 'decision' in text.lower() and 'uncertainty' in text.lower():
                # Check formatting of runs
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'neural correlates' in run.text.lower():
                            is_bold = run.font.bold is True
                            color = get_run_color_rgb(run)
                            is_white = color and color.upper() == 'FFFFFF'
                            size_ok = run.font.size is not None and abs(run.font.size - Pt(54)) < Pt(4)
                            if is_bold and is_white and size_ok:
                                title_found = True
                                print(f"  Title run: bold={is_bold}, color={color}, size={run.font.size}")
                            break
                    if title_found:
                        break
            if title_found:
                break
        if title_found:
            print(f"PASS: Component 2 -- Title text found with correct formatting (54pt, bold, white) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 -- Title not found with correct formatting (54pt bold white)")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Author line present, approximately 24pt, white (0.10 points)
    try:
        author_found = False
        for shape in all_shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            text = (shape.text or '').strip().lower()
            # Author line should contain names and possibly department/university info
            if any(kw in text for kw in ['dr.', 'prof.', 'department', 'university']):
                # Check it's not the title shape
                if 'neural correlates' not in text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                color = get_run_color_rgb(run)
                                is_white = color and color.upper() == 'FFFFFF'
                                size_ok = run.font.size is not None and abs(run.font.size - Pt(24)) < Pt(4)
                                if is_white and size_ok:
                                    author_found = True
                                    print(f"  Author run: color={color}, size={run.font.size}")
                                break
                        if author_found:
                            break
            if author_found:
                break
        if author_found:
            print(f"PASS: Component 3 -- Author line found (24pt, white) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 -- Author line not found with correct formatting")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Three-column layout with 7 required sections (0.20 points)
    # Required sections: Abstract, Introduction, Methods, Results, Discussion, Conclusions, References
    try:
        required_sections = ['abstract', 'introduction', 'methods', 'results', 'discussion', 'conclusions', 'references']
        found_sections = set()
        for shape in all_shapes:
            if not hasattr(shape, 'text'):
                continue
            text = (shape.text or '').strip().lower()
            for section in required_sections:
                # Section header should be a short text that matches the section name
                if text == section or text.startswith(section + '\n') or text.startswith(section + ' '):
                    # Also check this is a header-like shape (relatively short text)
                    if len(text) < 50:
                        found_sections.add(section)

        section_count = len(found_sections)
        missing = set(required_sections) - found_sections

        if section_count >= 7:
            print(f"PASS: Component 4 -- All 7 sections found: {found_sections} (0.20 pts)")
            total_score += 0.20
        elif section_count >= 5:
            partial = 0.20 * (section_count / 7)
            print(f"PARTIAL: Component 4 -- {section_count}/7 sections found, missing: {missing} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Only {section_count}/7 sections found: {found_sections}, missing: {missing}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Section headers are 28pt bold #003366 (0.15 points)
    try:
        correct_headers = 0
        for shape in all_shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            text = (shape.text or '').strip().lower()
            if text in required_sections:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip().lower() in required_sections:
                            is_bold = run.font.bold is True
                            color = get_run_color_rgb(run)
                            is_correct_color = color and color.upper() == '003366'
                            size_ok = run.font.size is not None and abs(run.font.size - Pt(28)) < Pt(4)
                            if is_bold and is_correct_color and size_ok:
                                correct_headers += 1
                            break
                    break

        if correct_headers >= 7:
            print(f"PASS: Component 5 -- All 7 section headers properly formatted (28pt, bold, #003366) (0.15 pts)")
            total_score += 0.15
        elif correct_headers >= 4:
            partial = 0.15 * (correct_headers / 7)
            print(f"PARTIAL: Component 5 -- {correct_headers}/7 headers correctly formatted ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 -- Only {correct_headers}/7 headers correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Col 1 sections (Abstract, Introduction) have #E8F0FE background (0.10 points)
    try:
        e8f0fe_count = 0
        for shape in all_shapes:
            fill_rgb = get_shape_fill_rgb(shape)
            if fill_rgb and fill_rgb.upper() == 'E8F0FE':
                e8f0fe_count += 1

        if e8f0fe_count >= 2:
            print(f"PASS: Component 6 -- Found {e8f0fe_count} shapes with #E8F0FE background (Col 1) (0.10 pts)")
            total_score += 0.10
        elif e8f0fe_count >= 1:
            print(f"PARTIAL: Component 6 -- Found {e8f0fe_count} shape(s) with #E8F0FE (expected 2) (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 6 -- No shapes with #E8F0FE background found")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Placeholders for brain scan image and bar chart (0.10 points)
    try:
        brain_scan_found = False
        bar_chart_found = False
        for shape in all_shapes:
            text = (shape.text or '').lower()
            if 'brain scan' in text or 'fmri' in text:
                brain_scan_found = True
            if 'bar chart' in text or 'bold signal' in text:
                bar_chart_found = True

        placeholders_found = int(brain_scan_found) + int(bar_chart_found)
        if placeholders_found == 2:
            print(f"PASS: Component 7 -- Both brain scan and bar chart placeholders found (0.10 pts)")
            total_score += 0.10
        elif placeholders_found == 1:
            print(f"PARTIAL: Component 7 -- Only {placeholders_found}/2 placeholders found (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 7 -- No image/chart placeholders found")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: University logo placeholder in top-right area (0.05 points)
    try:
        logo_found = False
        slide_width = prs.slide_width
        for shape in all_shapes:
            text = (shape.text or '').lower()
            if 'logo' in text or 'university' in text.lower():
                # Check if it's in the top-right area (right half and top portion)
                if shape.left > slide_width / 2 and shape.top < Inches(6):
                    logo_found = True
                    break
        if logo_found:
            print(f"PASS: Component 8 -- University logo placeholder in top-right (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 -- No university logo placeholder in top-right area")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
