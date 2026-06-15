"""
Initial Setup: Create Financial_Summary.pptx with 8 slides, slides 4 and 5 with titles but empty content
Task ID: impress_gf4_004
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_body(slide, text, font_size=Pt(18), bold=False):
    """Helper: add text to the content placeholder (index 1) if it exists."""
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = text
        for run in tf.paragraphs[0].runs:
            run.font.size = font_size
            run.font.bold = bold


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Acme Corp Financial Summary"
    slide1.placeholders[1].text = "Fiscal Year 2025 Overview"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Acme Corp achieved record revenue in FY2025, driven by strong growth in the enterprise segment."
    p2 = tf2.add_paragraph()
    p2.text = "Total revenue reached $12.8M, representing a 23% year-over-year increase."
    p3 = tf2.add_paragraph()
    p3.text = "Operating margins improved to 18.5%, up from 15.2% in the prior year."
    p4 = tf2.add_paragraph()
    p4.text = "Employee headcount grew to 245, with key hires in engineering and sales."

    # --- Slide 3: Revenue Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Highlights"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Enterprise segment contributed 62% of total revenue ($7.9M)."
    p = tf3.add_paragraph()
    p.text = "SMB segment grew 31% to $3.2M, exceeding internal targets."
    p = tf3.add_paragraph()
    p.text = "International markets now represent 28% of total revenue."
    p = tf3.add_paragraph()
    p.text = "Customer retention rate improved to 94.7% from 91.3%."

    # --- Slide 4: Quarterly Breakdown (empty content area) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box since blank layout has no title placeholder
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = txBox.text_frame
    tf4.text = "Quarterly Breakdown"
    for run in tf4.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Content area intentionally left empty - agent should add table here

    # --- Slide 5: Profit Distribution (empty content area) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf5 = txBox5.text_frame
    tf5.text = "Profit Distribution"
    for run in tf5.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # Content area intentionally left empty - agent should add pie chart here

    # --- Slide 6: Cost Analysis ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Cost Analysis"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Total operating costs were $10.4M, up 18% from FY2024."
    p = tf6.add_paragraph()
    p.text = "Personnel costs remain the largest expense at 58% of total costs."
    p = tf6.add_paragraph()
    p.text = "Infrastructure costs decreased 7% due to cloud migration savings."
    p = tf6.add_paragraph()
    p.text = "R&D investment increased to 22% of revenue, supporting product roadmap."

    # --- Slide 7: Regional Performance ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Regional Performance"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "North America: $9.2M revenue (72% of total), 21% growth."
    p = tf7.add_paragraph()
    p.text = "Europe: $2.3M revenue (18% of total), 28% growth."
    p = tf7.add_paragraph()
    p.text = "Asia-Pacific: $1.3M revenue (10% of total), 42% growth."
    p = tf7.add_paragraph()
    p.text = "LATAM pilot launched in Q3 with promising early traction."

    # --- Slide 8: Outlook & Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Outlook & Next Steps"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "FY2026 revenue target: $16.0M (25% growth)."
    p = tf8.add_paragraph()
    p.text = "Key initiatives: Enterprise platform v2.0 launch, LATAM expansion."
    p = tf8.add_paragraph()
    p.text = "Planned headcount increase to 310 employees by year-end."
    p = tf8.add_paragraph()
    p.text = "Board approval sought for Series C funding round."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
