"""
Reward Script: Verify table on slide 4, pie chart on slide 5, and hyperlink from Q1 to slide 5
Task ID: impress_gf4_004
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Table on slide 4 with 5 rows x 4 columns
  Component 2 (0.20): Header row has dark background, white bold text, correct column names
  Component 3 (0.20): Data rows contain Q1-Q4 with numeric financial data
  Component 4 (0.20): Pie chart on slide 5 with 4 profit segments (Q1-Q4)
  Component 5 (0.20): Hyperlink on Q1 cell in slide 4 table navigating to slide 5
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_004'


def persist_app_state(domain: str):
    """Save any unsaved changes in LibreOffice before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 5:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]
    slide5 = prs.slides[4]

    # Find table on slide 4
    table_shape = None
    for shape in slide4.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
            break

    # =========================================================================
    # Component 1: Table exists on slide 4 with correct dimensions (0.20 pts)
    # =========================================================================
    try:
        if table_shape is not None:
            tbl = table_shape.table
            nrows = len(tbl.rows)
            ncols = len(tbl.columns)
            if nrows == 5 and ncols == 4:
                print(f"PASS: Component 1 — Table on slide 4 is 5x4 (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Table is {nrows}x{ncols}, expected 5x4")
        else:
            print("FAIL: Component 1 — No table found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if table_shape is None:
        # No table means remaining table-related checks will all fail
        print("FAIL: Components 2-3, 5 — No table to verify")
        # Still check chart on slide 5
        # Component 4: Pie chart
        try:
            chart_found = False
            for shape in slide5.shapes:
                try:
                    if shape.has_chart:
                        chart = shape.chart
                        chart_type_val = chart.chart_type
                        # PIE = 5
                        if chart_type_val == 5:
                            series = chart.series
                            if len(series) >= 1 and len(series[0].values) == 4:
                                print(f"PASS: Component 4 — Pie chart on slide 5 with 4 segments (0.20 pts)")
                                total_score += 0.20
                                chart_found = True
                except:
                    pass
            if not chart_found:
                print("FAIL: Component 4 — No pie chart with 4 segments found on slide 5")
        except Exception as e:
            print(f"ERROR: Component 4 — {e}")

        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    tbl = table_shape.table

    # =========================================================================
    # Component 2: Header row dark background, white bold text, correct names (0.20 pts)
    # =========================================================================
    try:
        expected_headers = ['Quarter', 'Revenue', 'Costs', 'Profit']
        header_checks_passed = 0
        total_header_checks = 4  # one per column

        for c in range(4):
            cell = tbl.cell(0, c)
            cell_text = cell.text.strip()

            # Check header text (case-insensitive)
            if cell_text.lower() != expected_headers[c].lower():
                print(f"  FAIL: Header col {c} text is {repr(cell_text)}, expected {repr(expected_headers[c])}")
                continue

            # Check bold
            is_bold = False
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    if run.font.bold is True:
                        is_bold = True
                        break

            # Check dark fill on cell
            has_dark_fill = False
            try:
                fill = cell.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    rgb = fill.fore_color.rgb
                    # Dark means low RGB values (< 128 average)
                    r, g, b = rgb[0], rgb[1], rgb[2]
                    if isinstance(r, int):
                        avg = (r + g + b) / 3
                    else:
                        # RGBColor behaves like string 'RRGGBB'
                        hex_str = str(rgb)
                        r = int(hex_str[0:2], 16)
                        g = int(hex_str[2:4], 16)
                        b = int(hex_str[4:6], 16)
                        avg = (r + g + b) / 3
                    has_dark_fill = avg < 128
            except:
                pass

            # Check white text color
            has_white_text = False
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb)
                            r = int(rgb_str[0:2], 16)
                            g = int(rgb_str[2:4], 16)
                            b = int(rgb_str[4:6], 16)
                            if r > 200 and g > 200 and b > 200:
                                has_white_text = True
                    except:
                        pass

            if is_bold and has_dark_fill and has_white_text:
                header_checks_passed += 1
            else:
                print(f"  PARTIAL: Header col {c} ({cell_text}): bold={is_bold}, dark_fill={has_dark_fill}, white_text={has_white_text}")

        if header_checks_passed == 4:
            print(f"PASS: Component 2 — All 4 header cells have correct text, dark fill, white bold text (0.20 pts)")
            total_score += 0.20
        elif header_checks_passed >= 2:
            partial = 0.10
            print(f"PARTIAL: Component 2 — {header_checks_passed}/4 header cells correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {header_checks_passed}/4 header cells correct")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Data rows contain Q1-Q4 with financial data (0.20 pts)
    # =========================================================================
    try:
        expected_quarters = ['Q1', 'Q2', 'Q3', 'Q4']
        data_rows_ok = 0

        for r in range(1, 5):
            quarter_text = tbl.cell(r, 0).text.strip()
            if quarter_text != expected_quarters[r - 1]:
                print(f"  FAIL: Row {r} quarter is {repr(quarter_text)}, expected {repr(expected_quarters[r-1])}")
                continue

            # Check that revenue, costs, profit cells have numeric content
            has_numeric_data = True
            for c in range(1, 4):
                cell_text = tbl.cell(r, c).text.strip()
                # Remove common formatting: $, commas, spaces
                cleaned = cell_text.replace('$', '').replace(',', '').replace(' ', '').replace('.', '')
                if not cleaned.isdigit() or int(cleaned) == 0:
                    has_numeric_data = False
                    print(f"  FAIL: Row {r} col {c} has non-numeric or zero value: {repr(cell_text)}")
                    break

            if has_numeric_data:
                data_rows_ok += 1

        if data_rows_ok == 4:
            print(f"PASS: Component 3 — All 4 data rows (Q1-Q4) have valid financial data (0.20 pts)")
            total_score += 0.20
        elif data_rows_ok >= 2:
            partial = 0.10
            print(f"PARTIAL: Component 3 — {data_rows_ok}/4 data rows correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {data_rows_ok}/4 data rows correct")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Pie chart on slide 5 with 4 profit segments (0.20 pts)
    # =========================================================================
    try:
        chart_found = False
        for shape in slide5.shapes:
            try:
                if shape.has_chart:
                    chart = shape.chart
                    chart_type_val = chart.chart_type
                    # PIE = 5
                    is_pie = (chart_type_val == 5)
                    if not is_pie:
                        print(f"  INFO: Chart found but type is {chart_type_val}, expected PIE (5)")
                        continue

                    series = chart.series
                    if len(series) >= 1:
                        values = list(series[0].values)
                        num_points = len(values)

                        # Check categories for Q1-Q4
                        try:
                            cats = [str(c) for c in chart.plots[0].categories]
                        except:
                            cats = []

                        has_quarter_cats = all(
                            q in cats for q in ['Q1', 'Q2', 'Q3', 'Q4']
                        ) if len(cats) == 4 else False

                        if num_points == 4 and has_quarter_cats:
                            print(f"PASS: Component 4 — Pie chart with 4 Q1-Q4 segments, values={values} (0.20 pts)")
                            total_score += 0.20
                            chart_found = True
                        elif num_points == 4:
                            # Pie chart with 4 segments but without Q1-Q4 labels - partial
                            print(f"PARTIAL: Component 4 — Pie chart has 4 segments but categories are {cats} (0.10 pts)")
                            total_score += 0.10
                            chart_found = True
                        else:
                            print(f"FAIL: Component 4 — Pie chart has {num_points} segments, expected 4")
            except:
                pass

        if not chart_found:
            print("FAIL: Component 4 — No pie chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Hyperlink on Q1 cell navigating to slide 5 (0.20 pts)
    # =========================================================================
    try:
        q1_cell = tbl.cell(1, 0)
        hyperlink_found = False

        for p in q1_cell.text_frame.paragraphs:
            for run in p.runs:
                # Check for hyperlink via XML - look for hlinkClick with ppaction://hlinksldjump
                from lxml import etree
                r_elem = run._r
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

                hlink = r_elem.find('.//a:hlinkClick', ns)
                if hlink is not None:
                    action = hlink.get('action', '')
                    if 'hlinksldjump' in action:
                        hyperlink_found = True
                        print(f"PASS: Component 5 — Q1 cell has hyperlink with slide jump action (0.20 pts)")
                        total_score += 0.20
                        break
                    else:
                        # Has a hyperlink but not internal slide navigation
                        # Check if address points to slide 5
                        addr = run.hyperlink.address
                        if addr and ('slide5' in addr.lower() or '5' in addr):
                            hyperlink_found = True
                            print(f"PASS: Component 5 — Q1 cell has hyperlink to slide 5 (addr={addr}) (0.20 pts)")
                            total_score += 0.20
                            break
            if hyperlink_found:
                break

        if not hyperlink_found:
            # Also try checking via the hyperlink API directly
            try:
                for p in q1_cell.text_frame.paragraphs:
                    for run in p.runs:
                        if run.hyperlink and run.hyperlink.address:
                            addr = run.hyperlink.address
                            if 'slide' in addr.lower():
                                hyperlink_found = True
                                print(f"PASS: Component 5 — Q1 cell has hyperlink (addr={addr}) (0.20 pts)")
                                total_score += 0.20
                                break
                    if hyperlink_found:
                        break
            except:
                pass

        if not hyperlink_found:
            print("FAIL: Component 5 — No hyperlink found on Q1 cell in slide 4 table")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
