"""
Initial Setup: Create 8-slide presentation with transitions for PDF export task
Task ID: impress_el_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_033'
PPTX_PATH = f'{WORKDIR}/{TASK_ID}.pptx'
ODP_PATH = f'{WORKDIR}/Animated_Deck.odp'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Review"
    slide1.placeholders[1].text = "Presented by the Executive Strategy Team\nOctober 15, 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Analysis & Competitive Landscape"
    items = [
        "Revenue Performance by Region",
        "Product Development Pipeline",
        "Customer Acquisition & Retention Metrics",
        "Operational Efficiency Initiatives",
        "Strategic Partnerships Update",
        "Q4 Targets & Action Items",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Market Analysis ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Addressable Market: $42.8B (up 12% YoY)"
    data_points = [
        "North America: $18.2B market share, 23% penetration",
        "EMEA: $14.6B market share, 17% penetration",
        "Asia-Pacific: $7.3B, fastest growth at 28% CAGR",
        "Key competitor Nexus Corp launched rival platform in August",
        "Our differentiation: AI-powered analytics & enterprise integrations",
    ]
    for dp in data_points:
        p = body3.add_paragraph()
        p.text = dp
        p.level = 0

    # --- Slide 4: Revenue by Region ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Performance by Region"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Add table
    table_shape = slide4.shapes.add_table(
        5, 4, Inches(1), Inches(1.5), Inches(10), Inches(3)
    )
    table = table_shape.table
    headers = ["Region", "Q3 Revenue", "Q2 Revenue", "Growth %"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    rows_data = [
        ["North America", "$12.4M", "$11.1M", "+11.7%"],
        ["EMEA", "$8.7M", "$7.9M", "+10.1%"],
        ["Asia-Pacific", "$5.2M", "$4.1M", "+26.8%"],
        ["Latin America", "$2.1M", "$1.8M", "+16.7%"],
    ]
    for r, row_data in enumerate(rows_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Product Pipeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Product Development Pipeline"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Atlas 3.0 - Enterprise Analytics Suite (GA: Nov 2025)"
    features = [
        "Beacon - Real-time Monitoring Dashboard (Beta: Dec 2025)",
        "Compass - AI-driven Recommendation Engine (Alpha: Q1 2026)",
        "Drift Detection - Anomaly Alerts (GA: Oct 2025)",
        "Mobile SDK v2 - Cross-platform Support (GA: Nov 2025)",
    ]
    for f in features:
        p = body5.add_paragraph()
        p.text = f
        p.level = 0

    # --- Slide 6: Customer Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf6 = txBox6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Customer Acquisition & Retention"
    p6.font.size = Pt(28)
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    metrics = [
        ("New Enterprise Clients", "47", "+18% vs Q2"),
        ("Net Revenue Retention", "118%", "+3pp vs Q2"),
        ("Churn Rate", "2.1%", "-0.4pp vs Q2"),
        ("NPS Score", "72", "+5 vs Q2"),
        ("Avg. Deal Size", "$284K", "+22% vs Q2"),
    ]
    y_start = Inches(1.5)
    for i, (label, value, change) in enumerate(metrics):
        tb = slide6.shapes.add_textbox(Inches(1), y_start + Inches(i * 1.0), Inches(10), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}: "
        run.font.size = Pt(18)
        run.font.bold = True
        run2 = p.add_run()
        run2.text = f"{value}  "
        run2.font.size = Pt(24)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(0x00, 0x70, 0xC0)
        run3 = p.add_run()
        run3.text = change
        run3.font.size = Pt(14)
        run3.font.color.rgb = RGBColor(0x00, 0xB0, 0x50)

    # --- Slide 7: Strategic Partnerships ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Strategic Partnerships"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Salesforce Integration - Deep CRM sync (Live Q3)"
    partnerships = [
        "AWS Marketplace - Listed as Advanced Partner (Live Q3)",
        "Snowflake - Native data connector (Beta Q4)",
        "Microsoft Teams - Embedded analytics (Development)",
        "Okta SSO - Enterprise authentication (Live Q2)",
    ]
    for part in partnerships:
        p = body7.add_paragraph()
        p.text = part
        p.level = 0

    # --- Slide 8: Q4 Targets ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Q4 2025 Targets & Next Steps"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Revenue target: $32M (+14% QoQ)"
    targets = [
        "Launch Atlas 3.0 to 50 pilot enterprise accounts",
        "Expand APAC sales team by 15 headcount",
        "Achieve SOC 2 Type II certification by December",
        "Close 3 strategic partnership integrations",
        "Reduce customer onboarding time to under 48 hours",
    ]
    for t in targets:
        p = body8.add_paragraph()
        p.text = t
        p.level = 0

    prs.save(PPTX_PATH)
    print(f"Presentation created: {PPTX_PATH}")


def add_transitions_to_pptx(pptx_path):
    """Inject slide transitions into the PPTX XML."""
    # Transition types to apply to each slide (rotating through options)
    transitions = [
        '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>',
        '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:wipe dir="d"/></p:transition>',
        '<p:transition spd="slow" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:dissolve/></p:transition>',
        '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>',
        '<p:transition spd="fast" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:wipe dir="r"/></p:transition>',
        '<p:transition spd="slow" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:dissolve/></p:transition>',
        '<p:transition spd="med" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>',
        '<p:transition spd="fast" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:wipe dir="l"/></p:transition>',
    ]

    tmp_path = pptx_path + '.tmp'
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    ET.register_namespace('', 'http://schemas.openxmlformats.org/presentationml/2006/main')
    ET.register_namespace('a', 'http://schemas.openxmlformats.org/drawingml/2006/main')
    ET.register_namespace('r', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships')

    with zipfile.ZipFile(pptx_path, 'r') as zin, zipfile.ZipFile(tmp_path, 'w') as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                # Extract slide number
                slide_num = int(item.filename.replace('ppt/slides/slide', '').replace('.xml', ''))
                slide_idx = slide_num - 1
                if slide_idx < len(transitions):
                    # Parse XML, remove existing transitions, add new one
                    content = data.decode('utf-8')
                    root = ET.fromstring(content)
                    # Remove existing transition elements
                    for existing_tr in root.findall('.//p:transition', ns):
                        root.remove(existing_tr)
                    # Parse and append new transition
                    tr_elem = ET.fromstring(transitions[slide_idx])
                    root.append(tr_elem)
                    data = ET.tostring(root, encoding='unicode', xml_declaration=True).encode('utf-8')
            zout.writestr(item, data)

    os.replace(tmp_path, pptx_path)
    print(f"Transitions added to {pptx_path}")


def convert_to_odp_via_libreoffice(pptx_path, odp_path):
    """Use LibreOffice command line to convert PPTX to ODP."""
    # Kill any running LibreOffice first
    subprocess.run(['pkill', '-f', 'soffice'], capture_output=True)
    time.sleep(2)

    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    env["HOME"] = "/home/user"

    result = subprocess.run(
        [
            'libreoffice', '--headless', '--convert-to', 'odp',
            '--outdir', WORKDIR, pptx_path
        ],
        capture_output=True, text=True, env=env, timeout=60,
    )
    print(f"Conversion stdout: {result.stdout}")
    print(f"Conversion stderr: {result.stderr}")

    # The output will be named <task_id>.odp, rename to Animated_Deck.odp
    converted = os.path.join(WORKDIR, TASK_ID + '.odp')
    if os.path.exists(converted):
        if converted != odp_path:
            os.rename(converted, odp_path)
        print(f"Converted to: {odp_path}")
    else:
        print(f"WARNING: conversion output {converted} not found, listing dir:")
        print(os.listdir(WORKDIR))


def create_initial():
    # Step 1: Create PPTX with content
    create_presentation()

    # Step 2: Add transitions via XML
    add_transitions_to_pptx(PPTX_PATH)

    # Step 3: Convert to ODP
    convert_to_odp_via_libreoffice(PPTX_PATH, ODP_PATH)

    # Step 4: Verify ODP exists
    if os.path.exists(ODP_PATH):
        print(f"ODP file ready: {ODP_PATH} ({os.path.getsize(ODP_PATH)} bytes)")
    else:
        print("ERROR: ODP file was not created!")
        return

    # Step 5: Clean up the intermediate PPTX
    if os.path.exists(PPTX_PATH):
        os.remove(PPTX_PATH)

    # Step 6: Open in LibreOffice Impress (GUI ready)
    time.sleep(1)
    launch_gui(f'libreoffice --impress "{ODP_PATH}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
