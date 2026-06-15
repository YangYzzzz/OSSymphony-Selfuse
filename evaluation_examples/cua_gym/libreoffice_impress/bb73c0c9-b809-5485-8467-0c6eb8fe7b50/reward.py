"""
Reward Script: Import notes from docx files into LibreOffice Impress presentation
Task ID: osworld_multi_apps_impress_notes_import_010
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.4 pts): Slides 1-3 notes match intro_notes.docx content
  - Component 2 (0.4 pts): Slides 4-10 notes match main_notes.docx content
  - Component 3 (0.2 pts): All 10 slides have non-empty notes (complete coverage)
"""

import os

from pptx import Presentation
from docx import Document

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_import_010'

# Expected notes per slide (derived from the docx source files)
# intro_notes.docx provides notes for slides 1-3 (paragraphs at indices 2, 4, 6)
EXPECTED_NOTES = {
    1: 'Welcome the audience warmly. Introduce yourself and co-presenters. Mention that the session will be recorded and available for replay. Remind attendees to use the Q&A panel for questions throughout the webinar.',
    2: 'Walk through the agenda briefly \u2014 do not spend more than 2 minutes here. Highlight that the cost optimization section includes a live demo. Mention the case study at the end and encourage questions during the Q&A.',
    3: 'Explain cloud computing in simple terms for any non-technical attendees. Use the analogy: just like electricity, you pay for what you use. Emphasize the shift from capital expenditure to operational expenditure. Ask the audience if anyone is currently using cloud services \u2014 use polling feature.',
    4: 'Give a brief overview of each major provider. AWS has the largest market share (~32%). Azure is strong for enterprise Microsoft shops. GCP excels in data analytics and machine learning workloads. Mention that choice often depends on existing vendor relationships and specific services needed.',
    5: 'Distinguish clearly between IaaS, PaaS, and SaaS with concrete examples. IaaS example: EC2 virtual machines. PaaS example: Heroku or App Engine. SaaS example: Salesforce, Office 365. Emphasize that most organizations use a mix of all three models.',
    6: 'Explain when to choose each deployment model. Public cloud: startups and SMBs with variable workloads. Private cloud: regulated industries like banking and healthcare. Hybrid cloud: enterprises needing flexibility with some on-prem requirements. Note that multi-cloud strategies are increasingly common to avoid vendor lock-in.',
    7: 'This slide has the most actionable takeaways \u2014 slow down here. Right-sizing: audit resource usage every quarter, downsize underutilized instances. Reserved instances can save 30-60% vs on-demand for predictable workloads. Auto-scaling prevents over-provisioning during off-peak hours. Mention FinOps as an emerging discipline for cloud cost management.',
    8: 'Security is a shared responsibility model \u2014 cloud provider vs customer. Provider secures the infrastructure; customer secures data and access. Least privilege: grant only necessary permissions \u2014 use IAM roles, not root accounts. Compliance: SOC 2, ISO 27001, GDPR are common frameworks. Recommend regular penetration testing and vulnerability assessments.',
    9: 'Walk through the PayStream case study in detail \u2014 this is a real-world success story. Migration was done in 3 phases over 18 months to minimize business disruption. Phase 1: Non-critical workloads (dev/test environments). Phase 2: Data warehousing and analytics. Phase 3: Core banking applications (required most testing and compliance validation). Key lesson: involve security and compliance teams from day one.',
    10: 'Thank the audience again and acknowledge any specific questions from the Q&A panel. Reiterate the resource URL: cloud-webinar.example.com \u2014 resources available for 90 days. Mention the follow-up email survey \u2014 feedback helps improve future webinars. Next webinar topic: Kubernetes and Container Orchestration \u2014 registration link in chat. Stay on the call for an optional 10-minute extended Q&A for those interested.',
}


def get_slide_notes(slide):
    """Get notes text from a slide, returning empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def normalize_notes(text):
    """Normalize whitespace in notes text for comparison."""
    return ' '.join(text.split())


def verify_task(pptx_path):
    """
    Verify that notes from docx files were correctly imported into the presentation.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: ensure the file exists and can be loaded
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation {pptx_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify correct number of slides (should be 10)
    if len(prs.slides) != 10:
        print(f"CRITICAL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slides 1-3 notes match intro_notes.docx content (0.4 points)
    # Each slide in 1-3 is worth 0.4/3 ≈ 0.133 points
    try:
        intro_slides_correct = 0
        for slide_num in [1, 2, 3]:
            slide = prs.slides[slide_num - 1]
            actual_notes = get_slide_notes(slide)
            expected = EXPECTED_NOTES[slide_num]
            if normalize_notes(actual_notes) == normalize_notes(expected):
                print(f"PASS: Slide {slide_num} notes match intro_notes.docx content")
                intro_slides_correct += 1
            else:
                actual_preview = actual_notes[:60] if actual_notes else "(empty)"
                expected_preview = expected[:60]
                print(f"FAIL: Slide {slide_num} notes mismatch — got: '{actual_preview}...', expected: '{expected_preview}...'")

        intro_score = round(intro_slides_correct / 3 * 0.4, 4)
        if intro_slides_correct > 0:
            total_score += intro_score
        print(f"Component 1: {intro_slides_correct}/3 intro slides correct → +{intro_score} pts")
    except Exception as e:
        print(f"ERROR: Component 1 (intro slides check) — {e}")

    # Component 2: Slides 4-10 notes match main_notes.docx content (0.4 points)
    # Each slide in 4-10 is worth 0.4/7 ≈ 0.057 points
    try:
        main_slides_correct = 0
        for slide_num in range(4, 11):
            slide = prs.slides[slide_num - 1]
            actual_notes = get_slide_notes(slide)
            expected = EXPECTED_NOTES[slide_num]
            if normalize_notes(actual_notes) == normalize_notes(expected):
                print(f"PASS: Slide {slide_num} notes match main_notes.docx content")
                main_slides_correct += 1
            else:
                actual_preview = actual_notes[:60] if actual_notes else "(empty)"
                expected_preview = expected[:60]
                print(f"FAIL: Slide {slide_num} notes mismatch — got: '{actual_preview}...', expected: '{expected_preview}...'")

        main_score = round(main_slides_correct / 7 * 0.4, 4)
        if main_slides_correct > 0:
            total_score += main_score
        print(f"Component 2: {main_slides_correct}/7 main slides correct → +{main_score} pts")
    except Exception as e:
        print(f"ERROR: Component 2 (main slides check) — {e}")

    # Component 3: All 10 slides have non-empty notes (complete coverage) (0.2 points)
    # This verifies comprehensive note importation — all slides were addressed
    try:
        slides_with_notes = 0
        for i, slide in enumerate(prs.slides):
            notes = get_slide_notes(slide)
            if notes:
                slides_with_notes += 1
            else:
                print(f"FAIL: Slide {i+1} has empty notes")

        if slides_with_notes == 10:
            print(f"PASS: All 10 slides have non-empty notes")
            total_score += 0.2
            print(f"Component 3: All slides have notes → +0.2 pts")
        else:
            print(f"FAIL: Component 3 — only {slides_with_notes}/10 slides have notes")
    except Exception as e:
        print(f"ERROR: Component 3 (all slides coverage) — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM
file_path = f'{WORKDIR}/Webinar_Deck.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
