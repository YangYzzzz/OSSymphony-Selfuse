"""
Initial Setup: Webinar slides with empty notes, plus two notes docx files on Desktop
Task ID: osworld_multi_apps_impress_notes_import_010
Domain: libreoffice_impress (multi-app: impress + writer)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from docx import Document

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_import_010'
PPTX_OUTPUT = f'{WORKDIR}/Webinar_Deck.pptx'
INTRO_NOTES = f'{DESKTOP}/intro_notes.docx'
MAIN_NOTES  = f'{DESKTOP}/main_notes.docx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_webinar_deck():
    """Create Webinar_Deck.pptx with 10 slides and empty notes."""
    prs = Presentation()

    # Use a standard widescreen template size
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide content: realistic webinar slides
    slides_content = [
        # (title, body_text)
        ("Introduction to Cloud Computing", "Welcome to our webinar on modern cloud infrastructure."),
        ("Agenda", "1. Cloud Fundamentals\n2. Key Providers\n3. Cost Optimization\n4. Security Best Practices"),
        ("What is Cloud Computing?", "Cloud computing delivers computing services over the internet:\n- Servers\n- Storage\n- Databases\n- Networking"),
        ("Major Cloud Providers", "Amazon Web Services (AWS)\nMicrosoft Azure\nGoogle Cloud Platform (GCP)\nAlibaba Cloud"),
        ("Cloud Service Models", "IaaS – Infrastructure as a Service\nPaaS – Platform as a Service\nSaaS – Software as a Service"),
        ("Deployment Models", "Public Cloud: Shared infrastructure managed by provider\nPrivate Cloud: Dedicated infrastructure\nHybrid Cloud: Combination of public and private"),
        ("Cost Optimization Strategies", "Right-sizing instances\nReserved vs On-demand pricing\nAuto-scaling policies\nSpot instances for batch workloads"),
        ("Security Best Practices", "Principle of least privilege\nEncryption at rest and in transit\nMulti-factor authentication\nRegular audits and compliance checks"),
        ("Case Study: FinTech Migration", "Company: PayStream Financial\nChallenge: Legacy on-premise infrastructure\nSolution: Phased migration to AWS\nResult: 40% cost reduction, 99.99% uptime"),
        ("Q&A and Next Steps", "Thank you for attending!\nResources available at: cloud-webinar.example.com\nContact: webinar@techsolutions.com"),
    ]

    layout_title = prs.slide_layouts[0]   # Title Slide
    layout_content = prs.slide_layouts[1]  # Title + Content

    for i, (title_text, body_text) in enumerate(slides_content):
        if i == 0:
            slide = prs.slides.add_slide(layout_title)
            slide.shapes.title.text = title_text
            # subtitle placeholder
            try:
                slide.placeholders[1].text = body_text
            except (KeyError, IndexError):
                pass
        else:
            slide = prs.slides.add_slide(layout_content)
            slide.shapes.title.text = title_text
            try:
                slide.placeholders[1].text = body_text
            except (KeyError, IndexError):
                pass
        # NOTE: Do NOT set any notes — all slides must have empty notes in initial state

    prs.save(PPTX_OUTPUT)
    print(f'Initial file created: {PPTX_OUTPUT}')


def create_intro_notes():
    """Create intro_notes.docx with notes for slides 1-3."""
    doc = Document()

    doc.add_heading('Webinar Introduction Notes', level=1)

    doc.add_heading('Slide 1 Notes', level=2)
    doc.add_paragraph(
        'Welcome the audience warmly. Introduce yourself and co-presenters. '
        'Mention that the session will be recorded and available for replay. '
        'Remind attendees to use the Q&A panel for questions throughout the webinar.'
    )

    doc.add_heading('Slide 2 Notes', level=2)
    doc.add_paragraph(
        'Walk through the agenda briefly — do not spend more than 2 minutes here. '
        'Highlight that the cost optimization section includes a live demo. '
        'Mention the case study at the end and encourage questions during the Q&A.'
    )

    doc.add_heading('Slide 3 Notes', level=2)
    doc.add_paragraph(
        'Explain cloud computing in simple terms for any non-technical attendees. '
        'Use the analogy: just like electricity, you pay for what you use. '
        'Emphasize the shift from capital expenditure to operational expenditure. '
        'Ask the audience if anyone is currently using cloud services — use polling feature.'
    )

    doc.save(INTRO_NOTES)
    print(f'Intro notes created: {INTRO_NOTES}')


def create_main_notes():
    """Create main_notes.docx with notes for slides 4-10."""
    doc = Document()

    doc.add_heading('Webinar Main Session Notes', level=1)

    notes_data = [
        (
            'Slide 4 Notes',
            'Give a brief overview of each major provider. '
            'AWS has the largest market share (~32%). Azure is strong for enterprise Microsoft shops. '
            'GCP excels in data analytics and machine learning workloads. '
            'Mention that choice often depends on existing vendor relationships and specific services needed.'
        ),
        (
            'Slide 5 Notes',
            'Distinguish clearly between IaaS, PaaS, and SaaS with concrete examples. '
            'IaaS example: EC2 virtual machines. PaaS example: Heroku or App Engine. '
            'SaaS example: Salesforce, Office 365. '
            'Emphasize that most organizations use a mix of all three models.'
        ),
        (
            'Slide 6 Notes',
            'Explain when to choose each deployment model. '
            'Public cloud: startups and SMBs with variable workloads. '
            'Private cloud: regulated industries like banking and healthcare. '
            'Hybrid cloud: enterprises needing flexibility with some on-prem requirements. '
            'Note that multi-cloud strategies are increasingly common to avoid vendor lock-in.'
        ),
        (
            'Slide 7 Notes',
            'This slide has the most actionable takeaways — slow down here. '
            'Right-sizing: audit resource usage every quarter, downsize underutilized instances. '
            'Reserved instances can save 30-60% vs on-demand for predictable workloads. '
            'Auto-scaling prevents over-provisioning during off-peak hours. '
            'Mention FinOps as an emerging discipline for cloud cost management.'
        ),
        (
            'Slide 8 Notes',
            'Security is a shared responsibility model — cloud provider vs customer. '
            'Provider secures the infrastructure; customer secures data and access. '
            'Least privilege: grant only necessary permissions — use IAM roles, not root accounts. '
            'Compliance: SOC 2, ISO 27001, GDPR are common frameworks. '
            'Recommend regular penetration testing and vulnerability assessments.'
        ),
        (
            'Slide 9 Notes',
            'Walk through the PayStream case study in detail — this is a real-world success story. '
            'Migration was done in 3 phases over 18 months to minimize business disruption. '
            'Phase 1: Non-critical workloads (dev/test environments). '
            'Phase 2: Data warehousing and analytics. '
            'Phase 3: Core banking applications (required most testing and compliance validation). '
            'Key lesson: involve security and compliance teams from day one.'
        ),
        (
            'Slide 10 Notes',
            'Thank the audience again and acknowledge any specific questions from the Q&A panel. '
            'Reiterate the resource URL: cloud-webinar.example.com — resources available for 90 days. '
            'Mention the follow-up email survey — feedback helps improve future webinars. '
            'Next webinar topic: Kubernetes and Container Orchestration — registration link in chat. '
            'Stay on the call for an optional 10-minute extended Q&A for those interested.'
        ),
    ]

    for heading, content in notes_data:
        doc.add_heading(heading, level=2)
        doc.add_paragraph(content)

    doc.save(MAIN_NOTES)
    print(f'Main notes created: {MAIN_NOTES}')


def create_initial():
    # Ensure Desktop directory exists
    os.makedirs(DESKTOP, exist_ok=True)

    create_webinar_deck()
    create_intro_notes()
    create_main_notes()

    # GUI-ready startup: open Webinar_Deck.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{PPTX_OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with Webinar_Deck.pptx (DISPLAY=:0)')


create_initial()
