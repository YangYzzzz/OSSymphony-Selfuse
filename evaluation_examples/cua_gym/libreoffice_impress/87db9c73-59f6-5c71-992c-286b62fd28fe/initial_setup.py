"""
Initial Setup: Create exam instructions presentation with 5 slides.
Task ID: impress_tm_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and add body text to a slide."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Final Exam Instructions"
    slide1.placeholders[1].text = "Spring 2026 - Introduction to Computer Science\nProfessor Elena Marchetti"

    # --- Slide 2: General Guidelines ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "General Guidelines", [
        "You have 90 minutes to complete this exam",
        "Write your full name and student ID on every answer sheet",
        "Answer all questions in the order they appear",
        "Partial credit will be awarded for incomplete solutions",
        "Show all your work for calculation-based problems",
        "No electronic devices except approved calculators",
    ])

    # --- Slide 3: Allowed Materials ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Allowed Materials", [
        "One double-sided handwritten cheat sheet (8.5 x 11 inches)",
        "Non-programmable scientific calculator",
        "Standard writing instruments (pens, pencils, erasers)",
        "Blank scratch paper will be provided by the proctor",
        "No textbooks, printed notes, or digital references",
        "Water bottles with labels removed are permitted",
    ])

    # --- Slide 4: Read Carefully (target slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Read Carefully", [
        "IMPORTANT: Do NOT open the exam booklet until instructed",
        "Raise your hand if you have a question during the exam",
        "If you finish early, review your answers before submitting",
        "All exam materials must be returned to the proctor",
        "Academic integrity violations will result in automatic failure",
        "The exam starts when the proctor announces 'Begin'",
    ])

    # --- Slide 5: Good Luck ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Good Luck!", [
        "Take a deep breath and stay calm",
        "Read each question thoroughly before answering",
        "Budget your time - don't spend too long on any single question",
        "Trust your preparation and do your best",
        "Results will be posted within 5 business days",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
