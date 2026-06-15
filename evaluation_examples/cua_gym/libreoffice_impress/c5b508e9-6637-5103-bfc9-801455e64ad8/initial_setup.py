"""
Initial Setup: Create a presentation with 8 slides; slide 5 has white background with a chart.
Task ID: impress_el_080
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_el_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_chart_slide(prs, title_text):
    """Slide 5: white background with a bar chart showing quarterly revenue."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Explicit white background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # Add chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1 2025', 'Q2 2025', 'Q3 2025', 'Q4 2025']
    chart_data.add_series('Revenue ($M)', (12.4, 15.8, 14.2, 18.6))
    chart_data.add_series('Expenses ($M)', (9.1, 10.3, 11.0, 12.5))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(1.5), Inches(8.0), Inches(5.0),
        chart_data
    )

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Greenfield Technologies", "Annual Performance Review 2025")

    # Slide 2: Company Overview
    add_content_slide(prs, "Company Overview", [
        "Founded in 2018, headquartered in Austin, TX",
        "340 employees across 5 global offices",
        "Core focus: renewable energy solutions",
        "Serving enterprise clients in 28 countries",
    ])

    # Slide 3: Key Achievements
    add_content_slide(prs, "Key Achievements", [
        "Revenue growth of 42% year-over-year",
        "Launched 3 new product lines in Q2",
        "Customer satisfaction score: 94.7%",
        "Expanded to European and Asian markets",
        "Won 'Best Green Tech Startup' award at CES 2025",
    ])

    # Slide 4: Strategic Priorities
    add_content_slide(prs, "Strategic Priorities for 2026", [
        "Accelerate R&D investment in solar storage",
        "Hire 120 new engineers by Q3 2026",
        "Establish partnerships with utility providers",
        "ISO 14001 environmental certification",
    ])

    # Slide 5: Financial Overview with chart (white background)
    add_chart_slide(prs, "Financial Overview — Quarterly Performance")

    # Slide 6: Team Structure
    add_content_slide(prs, "Team Structure", [
        "Engineering: 145 staff (Director: Priya Sharma)",
        "Sales & Marketing: 82 staff (VP: David Reyes)",
        "Operations: 63 staff (COO: Elena Vasquez)",
        "Research: 50 staff (Chief Scientist: Dr. James Park)",
    ])

    # Slide 7: Product Roadmap
    add_content_slide(prs, "Product Roadmap", [
        "SolarVault 3.0 — residential battery storage (June 2026)",
        "WindSync Pro — industrial turbine management (Sept 2026)",
        "EcoGrid Platform — smart grid integration (Dec 2026)",
        "GreenFleet — EV charging infrastructure (March 2027)",
    ])

    # Slide 8: Thank You / Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.size = Pt(24)
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
