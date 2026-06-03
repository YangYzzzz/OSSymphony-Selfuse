"""
Reward Script: Duplicate slide 5 three times with colored backgrounds
Task ID: impress_el_080
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.20): Total slide count is 11
  - Component 2 (0.15): Slide 5 background unchanged (#FFFFFF)
  - Component 3 (0.20): Slide 6 background is #CC0000 (red copy)
  - Component 4 (0.20): Slide 7 background is #FFCC00 (yellow copy)
  - Component 5 (0.20): Slide 8 background is #00CC00 (green copy)
  - Component 6 (0.05): All three copies (slides 6-8) contain a chart
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_el_080'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_bg_color(slide):
    """Return background color as uppercase hex string (e.g. 'CC0000'), or None."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 11 (0.20 points)
    # Initial has 8 slides. Task adds 3 duplicates -> 11.
    try:
        if num_slides == 11:
            print(f"PASS: Component 1 — Slide count is 11 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Expected 11 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 8 slides to check the rest
    if num_slides < 8:
        print(f"CRITICAL: Not enough slides ({num_slides}) to verify backgrounds")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 5 background is still #FFFFFF (0.15 points)
    # This is scored only when slide count changed (otherwise initial would pass too)
    try:
        slide5_bg = get_slide_bg_color(prs.slides[4])  # 0-indexed
        if num_slides > 8 and slide5_bg == "FFFFFF":
            print(f"PASS: Component 2 — Slide 5 bg is #FFFFFF (unchanged) (0.15 pts)")
            total_score += 0.15
        elif num_slides <= 8:
            print(f"FAIL: Component 2 — Slide count not changed, skipping (no credit)")
        else:
            print(f"FAIL: Component 2 — Slide 5 bg expected #FFFFFF, found {slide5_bg}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 6 background is #CC0000 (0.20 points)
    try:
        if num_slides >= 9:
            slide6_bg = get_slide_bg_color(prs.slides[5])
            if slide6_bg == "CC0000":
                print(f"PASS: Component 3 — Slide 6 bg is #CC0000 (red) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 — Slide 6 bg expected #CC0000, found {slide6_bg}")
        else:
            print(f"FAIL: Component 3 — Slide 6 does not exist (only {num_slides} slides)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 7 background is #FFCC00 (0.20 points)
    try:
        if num_slides >= 10:
            slide7_bg = get_slide_bg_color(prs.slides[6])
            if slide7_bg == "FFCC00":
                print(f"PASS: Component 4 — Slide 7 bg is #FFCC00 (yellow) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 4 — Slide 7 bg expected #FFCC00, found {slide7_bg}")
        else:
            print(f"FAIL: Component 4 — Slide 7 does not exist (only {num_slides} slides)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 8 background is #00CC00 (0.20 points)
    try:
        if num_slides >= 11:
            slide8_bg = get_slide_bg_color(prs.slides[7])
            if slide8_bg == "00CC00":
                print(f"PASS: Component 5 — Slide 8 bg is #00CC00 (green) (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 5 — Slide 8 bg expected #00CC00, found {slide8_bg}")
        else:
            print(f"FAIL: Component 5 — Slide 8 does not exist (only {num_slides} slides)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slides 6-8 each contain a chart (copies retain content) (0.05 points)
    # Only award if all 3 copies have charts
    try:
        if num_slides >= 11:
            charts_found = 0
            for idx in [5, 6, 7]:  # slides 6, 7, 8 (0-indexed)
                slide = prs.slides[idx]
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        charts_found += 1
                        break
            if charts_found == 3:
                print(f"PASS: Component 6 — All 3 copies have charts (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 6 — Expected charts in all 3 copies, found {charts_found}/3")
        else:
            print(f"FAIL: Component 6 — Not enough slides to check chart copies")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification (LibreOffice may have unsaved changes)
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
