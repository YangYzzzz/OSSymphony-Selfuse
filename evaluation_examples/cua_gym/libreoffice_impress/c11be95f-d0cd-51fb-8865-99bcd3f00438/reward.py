"""
FINAL REWARD SCRIPT - SUCCESS
Task: Use Roman numerals (I, II, III) for pages 1–2, positioned top-center.
Generated: 2025-10-17 05:42:58
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation


def verify_roman_numerals_top_center(file_path: str) -> float:
    """
    Verify that the first two slides of a PPTX presentation have Roman numerals
    (I, II) positioned in the top-centre area of the slide.

    Scoring (progressive):
      • 0.45 pts for finding each correct numeral on its slide (I on slide-1, II on slide-2)
      • 0.05 pts for each numeral being within the required top-centre bounds
        (within top 15 % of slide height AND centre ±5 % of slide width)
      → Perfect score 1.0 when both slides contain the correct numeral and each is
        correctly positioned.
    """

    max_score = 1.0
    total_score = 0.0

    # ---------- 1. Load presentation ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not open PPTX: {e}")
        return 0.0

    slide_w, slide_h = prs.slide_width, prs.slide_height
    expected_numerals = ["I", "II"]  # slides 1 and 2 respectively

    # ---------- 2. Inspect each required slide ----------
    for slide_idx, numeral in enumerate(expected_numerals):
        if slide_idx >= len(prs.slides):
            print(f"✗ Slide {slide_idx+1} missing (expected numeral '{numeral}')")
            continue  # can't award anything for this slide

        slide = prs.slides[slide_idx]
        print(f"Checking slide {slide_idx+1} for numeral '{numeral}' ...")

        numeral_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() == numeral:
                numeral_shape = shape
                break

        # ----- 2a. Numeral correctness -----
        if numeral_shape is None:
            print(f"  ✗ Numeral '{numeral}' NOT found on slide {slide_idx+1}")
            continue  # cannot award position points either

        print("  ✓ Correct numeral found")
        total_score += 0.45  # award for correct numeral

        # ----- 2b. Position correctness -----
        try:
            left, top, width = numeral_shape.left, numeral_shape.top, numeral_shape.width
            centre_x = left + width / 2
            slide_centre_x = slide_w / 2

            horizontal_ok = abs(centre_x - slide_centre_x) <= slide_w * 0.05  # ±5 %
            vertical_ok = top <= slide_h * 0.15  # within top 15 %

            print(
                f"    Position left={left}, top={top}, width={width} | "
                f"centre offset={abs(centre_x - slide_centre_x)}"
            )
            if horizontal_ok and vertical_ok:
                print("    ✓ Position within top-centre bounds")
                total_score += 0.05
            else:
                print("    ✗ Numeral not in required top-centre area")
        except Exception as e:
            print(f"    ✗ Could not evaluate position: {e}")

    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task description
    presentation_path = "/home/user/use_roman_numerals_i_ii_iii_for_pages_12_positioned_top_center.pptx"
    reward = verify_roman_numerals_top_center(presentation_path)
    print(f"REWARD: {reward}")
