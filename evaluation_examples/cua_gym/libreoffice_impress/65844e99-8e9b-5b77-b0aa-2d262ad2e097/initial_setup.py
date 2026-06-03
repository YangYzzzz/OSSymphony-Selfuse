"""
Initial Setup: Export presentation as PDF with JPEG compression
Task ID: impstruct_047
Domain: libreoffice_impress

Creates a 15-slide presentation with high-resolution photographs
and opens it in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impstruct_047'
OUTPUT = f'{WORKDIR}/image_heavy.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def generate_photo_image(width, height, theme_color, label_text):
    """Generate a realistic-looking high-res photograph placeholder image."""
    img = Image.new('RGB', (width, height), theme_color)
    draw = ImageDraw.Draw(img)

    # Add gradient-like bands for visual interest
    r0, g0, b0 = theme_color
    for y in range(0, height, max(1, height // 20)):
        factor = 1.0 - 0.3 * (y / height)
        band_color = (
            max(0, min(255, int(r0 * factor))),
            max(0, min(255, int(g0 * factor))),
            max(0, min(255, int(b0 * factor))),
        )
        draw.rectangle([0, y, width, min(y + height // 20, height)], fill=band_color)

    # Add some geometric elements to simulate photo content
    cx, cy = width // 2, height // 2
    draw.ellipse([cx - 200, cy - 200, cx + 200, cy + 200],
                 fill=(min(255, r0 + 40), min(255, g0 + 40), min(255, b0 + 40)))

    # Add label text
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
    except (IOError, OSError):
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), label_text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((width - tw) // 2, (height - th) // 2), label_text,
              fill=(255, 255, 255), font=font)

    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


# Slide themes: title, color scheme, image label
slide_themes = [
    ("Annual Company Retreat 2025", (34, 85, 128), "Mountain Landscape"),
    ("Q1 Revenue Performance", (128, 45, 32), "Revenue Chart Photo"),
    ("New Product Launch Event", (45, 120, 78), "Product Showcase"),
    ("Team Building Activities", (90, 60, 130), "Team Activity Photo"),
    ("Global Market Expansion", (42, 100, 120), "World Map Overview"),
    ("Customer Satisfaction Survey Results", (150, 80, 30), "Survey Dashboard"),
    ("R&D Innovation Highlights", (60, 90, 140), "Laboratory Equipment"),
    ("Supply Chain Optimization", (100, 70, 50), "Warehouse Aerial View"),
    ("Employee Recognition Awards", (140, 50, 80), "Awards Ceremony"),
    ("Sustainability Initiatives", (30, 110, 60), "Solar Panel Installation"),
    ("Digital Transformation Roadmap", (70, 55, 130), "Technology Infrastructure"),
    ("Office Renovation Progress", (120, 95, 45), "Construction Progress"),
    ("Holiday Party Memories", (160, 40, 60), "Celebration Photo"),
    ("Safety Training Documentation", (50, 80, 100), "Training Session"),
    ("Year-End Financial Summary", (80, 70, 120), "Financial Overview"),
]


def create_initial():
    prs = Presentation()
    # Standard 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, (title, color, img_label) in enumerate(slide_themes):
        if i == 0:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if slide.placeholders[1]:
                slide.placeholders[1].text = "Presented by Corporate Communications"
        else:
            # Content slide with blank layout
            slide = prs.slides.add_slide(prs.slide_layouts[5])

            # Title text box at top
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # Generate and add a high-res image (1920x1080)
        img_buf = generate_photo_image(1920, 1080, color, img_label)

        if i == 0:
            # Smaller image on title slide
            slide.shapes.add_picture(
                img_buf, Inches(2), Inches(3.5), Inches(9), Inches(3.5)
            )
        else:
            # Large image filling most of the slide
            slide.shapes.add_picture(
                img_buf, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.8)
            )

        # Add a second smaller image to some slides for "many photographs" feel
        if i % 2 == 0 and i > 0:
            img_buf2 = generate_photo_image(
                1280, 720,
                (min(255, color[0] + 50), min(255, color[1] + 30), min(255, color[2] + 20)),
                f"Detail View {i}"
            )
            slide.shapes.add_picture(
                img_buf2, Inches(8.5), Inches(1.3), Inches(4.3), Inches(2.5)
            )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=3.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
