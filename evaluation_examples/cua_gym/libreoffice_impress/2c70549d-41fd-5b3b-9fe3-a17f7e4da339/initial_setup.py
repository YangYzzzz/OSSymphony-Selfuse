"""
Initial Setup: Create a 6-slide quarterly review presentation with no transitions.
Task ID: impress_tm_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(placeholder, text, font_name="Calibri", font_size=Pt(18),
                            bold=False, color=None, alignment=None):
    """Helper to set styled text on a placeholder."""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet_slide(slide, title_text, bullets, layout_name="content"):
    """Add title and bulleted content to a slide."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(16)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Quarterly Review"
    slide1.placeholders[1].text = "Prepared by Strategy & Operations Team\nMarch 31, 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide2, "Agenda", [
        "Financial performance overview",
        "Revenue breakdown by segment",
        "Team performance metrics",
        "Key project milestones",
        "Risks and mitigation strategies",
        "Next quarter priorities and action items",
    ])

    # --- Slide 3: Financial Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide3, "Financial Highlights", [
        "Total revenue: $12.4M (+8.2% YoY)",
        "Gross margin: 62.3% (target: 60%)",
        "Operating expenses: $4.1M (-3.5% vs Q4)",
        "Net income: $2.8M, exceeding forecast by $340K",
        "Customer acquisition cost: $127 (down from $145)",
    ])

    # --- Slide 4: Team Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide4, "Team Performance", [
        "Engineering: Shipped 14 features, 98.7% uptime",
        "Sales: 112% quota attainment, 23 new accounts",
        "Marketing: 45K qualified leads, 22% conversion rate",
        "Support: Average resolution time 4.2 hours (SLA: 8h)",
        "Employee satisfaction score: 4.3/5.0",
    ])

    # --- Slide 5: Key Projects ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide5, "Key Projects", [
        "Project Atlas: Platform migration 85% complete (on track)",
        "Project Beacon: New analytics dashboard launched Feb 15",
        "Project Catalyst: API v3 beta testing with 12 partners",
        "Project Delta: Mobile app redesign in UX review phase",
    ])

    # --- Slide 6: Next Steps & Q&A ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide6, "Next Steps & Q&A", [
        "Finalize Q2 OKRs by April 7",
        "Submit budget proposals by April 14",
        "Schedule cross-functional planning sessions",
        "Begin Project Atlas final migration phase",
        "Questions and open discussion",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
