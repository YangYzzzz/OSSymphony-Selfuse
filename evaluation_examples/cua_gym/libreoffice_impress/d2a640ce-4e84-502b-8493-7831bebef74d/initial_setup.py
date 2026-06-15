"""
Initial Setup: Create a 6-slide Regional Review presentation with slide 4 empty for chart insertion.
Task ID: impress_gf2_025
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_025'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet_slide(slide, title_text, bullets):
    """Add title and bulleted content to a slide."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(16)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Regional Annual Review 2023"
    slide1.placeholders[1].text = "Performance Analysis Across All Regions"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide2, "Executive Summary", [
        "Overall revenue grew 12.3% year-over-year reaching $48.7M",
        "North and West regions showed strongest growth trajectories",
        "East region experienced slight decline due to market saturation",
        "Central region exceeded expectations with 19% improvement",
        "Customer retention rate improved from 87% to 91% across all regions",
    ])

    # --- Slide 3: Revenue Highlights ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Title textbox
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue Highlights by Region"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Revenue table
    rows, cols = 6, 4
    table_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5))
    table = table_shape.table
    headers = ["Region", "2022 Revenue ($K)", "2023 Revenue ($K)", "Growth %"]
    data = [
        ["North", "120", "145", "+20.8%"],
        ["South", "95", "108", "+13.7%"],
        ["East", "140", "132", "-5.7%"],
        ["West", "88", "120", "+36.4%"],
        ["Central", "105", "125", "+19.0%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        from pptx.dml.color import RGBColor as RC
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x7A)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # --- Slide 4: Year-over-Year Regional Performance (EMPTY - for chart task) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title textbox only - no chart
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf4 = title_box4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Year-over-Year Regional Performance"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # --- Slide 5: Strategic Initiatives ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullet_slide(slide5, "Strategic Initiatives for 2024", [
        "Expand North region distribution network by Q2 2024",
        "Launch targeted marketing campaign in South region",
        "Restructure East region sales team to address decline",
        "Replicate West region success model in underperforming areas",
        "Invest in Central region infrastructure to sustain momentum",
    ])

    # --- Slide 6: Q&A ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Questions & Discussion"
    slide6.placeholders[1].text = "Thank you for your attention\nContact: regional-review@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
