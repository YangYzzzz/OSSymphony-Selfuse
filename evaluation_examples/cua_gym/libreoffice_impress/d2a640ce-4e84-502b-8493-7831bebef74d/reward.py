"""
Reward Script: Insert grouped bar chart on slide 4 with 2022 vs 2023 regional data
Task ID: impress_gf2_025
Domain: libreoffice_impress
Scoring:
  C1: Chart exists on slide 4, BAR_CLUSTERED type (0.15)
  C2: Chart title is 'Regional Performance Comparison' (0.10)
  C3: Two series named '2022' and '2023' (0.15)
  C4: Series '2022' values correct (0.15)
  C5: Series '2023' values correct (0.15)
  C6: Categories correct (0.10)
  C7: Data labels on 2023 series only (0.15)
  C8: Chart has legend (0.05)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_025'


def persist_app_state(domain: str):
    """Try to save any unsaved LibreOffice state."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        try:
            if shape.has_chart:
                chart_shape = shape
                break
        except Exception:
            pass

    # Component 1: Chart exists on slide 4 and is BAR_CLUSTERED (0.15 pts)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # BAR_CLUSTERED = 57
            if chart.chart_type == 57:
                print(f"PASS: Component 1 -- BAR_CLUSTERED chart found on slide 4 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 -- Chart exists but type is {chart.chart_type}, expected BAR_CLUSTERED (57)")
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 4")
            # No chart means no further checks can pass
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    chart = chart_shape.chart
    plot = chart.plots[0]

    # Component 2: Chart title is 'Regional Performance Comparison' (0.10 pts)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Regional Performance Comparison':
                print(f"PASS: Component 2 -- Chart title is correct (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Chart title is '{title_text}', expected 'Regional Performance Comparison'")
        else:
            print(f"FAIL: Component 2 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Two series named '2022' and '2023' (0.15 pts)
    try:
        num_series = len(plot.series)
        if num_series == 2:
            # Get series names via XML
            series_names = []
            for ser in plot.series:
                ser_el = ser._element
                tx = ser_el.find(qn('c:tx'), ser_el.nsmap)
                name = None
                if tx is not None:
                    v = tx.find('.//' + qn('c:v'), ser_el.nsmap)
                    if v is not None:
                        name = v.text
                series_names.append(name)

            if series_names == ['2022', '2023']:
                print(f"PASS: Component 3 -- Two series named '2022' and '2023' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- Series names are {series_names}, expected ['2022', '2023']")
        else:
            print(f"FAIL: Component 3 -- Found {num_series} series, expected 2")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Series '2022' values correct (0.15 pts)
    expected_2022 = [120.0, 95.0, 140.0, 88.0, 105.0]
    try:
        actual_2022 = list(plot.series[0].values)
        if actual_2022 == expected_2022:
            print(f"PASS: Component 4 -- 2022 series values correct: {actual_2022} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- 2022 series values are {actual_2022}, expected {expected_2022}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Series '2023' values correct (0.15 pts)
    expected_2023 = [145.0, 108.0, 132.0, 120.0, 125.0]
    try:
        actual_2023 = list(plot.series[1].values)
        if actual_2023 == expected_2023:
            print(f"PASS: Component 5 -- 2023 series values correct: {actual_2023} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 -- 2023 series values are {actual_2023}, expected {expected_2023}")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Categories are [North, South, East, West, Central] (0.10 pts)
    expected_cats = ['North', 'South', 'East', 'West', 'Central']
    try:
        actual_cats = list(plot.categories)
        if actual_cats == expected_cats:
            print(f"PASS: Component 6 -- Categories correct: {actual_cats} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 -- Categories are {actual_cats}, expected {expected_cats}")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Data labels on 2023 series ONLY (0.15 pts)
    # Series 0 (2022) should NOT have data labels, series 1 (2023) should have showVal=1
    try:
        ser0_el = plot.series[0]._element
        ser1_el = plot.series[1]._element

        def _has_show_val(ser_element):
            """Check if a series element has dLbls with showVal='1'."""
            dl = ser_element.find(qn('c:dLbls'), ser_element.nsmap)
            if dl is None:
                return False
            sv = dl.find(qn('c:showVal'), dl.nsmap)
            return sv is not None and sv.get('val') == '1'

        # Check series 0 has no data labels, series 1 has showVal=1
        series0_has_labels = _has_show_val(ser0_el)
        series1_has_labels = _has_show_val(ser1_el)

        if not series0_has_labels and series1_has_labels:
            print(f"PASS: Component 7 -- Data labels on 2023 series only (0.15 pts)")
            total_score += 0.15
        elif series1_has_labels and series0_has_labels:
            print(f"FAIL: Component 7 -- Data labels on BOTH series, should be 2023 only")
        elif not series1_has_labels:
            print(f"FAIL: Component 7 -- No data labels on 2023 series")
        else:
            print(f"FAIL: Component 7 -- Unexpected data label state: s0={series0_has_labels}, s1={series1_has_labels}")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Chart has legend (0.05 pts)
    try:
        if chart.has_legend:
            print(f"PASS: Component 8 -- Chart has legend (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 -- Chart has no legend")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
