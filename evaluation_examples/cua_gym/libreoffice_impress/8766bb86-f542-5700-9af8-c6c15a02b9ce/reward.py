"""
Reward Script: Add decorative rectangle bar at bottom of master slide
Task ID: impress_ma_020
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Non-placeholder rectangle shape exists on master slide
  Component 2 (0.4): Correct position and dimensions (full width, 0.3in tall, at bottom)
  Component 3 (0.3): Correct solid fill color #FF6B35
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_020'


def is_approx_equal(val1, val2, tolerance=0.02):
    """Check if two values are approximately equal within relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 and val2 == 0:
        return True
    denom = max(abs(val1), abs(val2))
    if denom == 0:
        return True
    return abs(val1 - val2) / denom <= tolerance


def find_decorative_rect(master):
    """Find non-placeholder rectangle shapes on the master slide."""
    rects = []
    for shape in master.shapes:
        # Skip placeholders -- they are pre-existing
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        # Look for auto shapes (rectangles)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            rects.append(shape)
    return rects


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide dimensions for reference
    slide_width = prs.slide_width   # expected: 9144000 EMU = 10 inches
    slide_height = prs.slide_height  # expected: 6858000 EMU = 7.5 inches

    # Get master slide
    if len(prs.slide_masters) == 0:
        print("FAIL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]

    # Find non-placeholder shapes on master
    decorative_shapes = find_decorative_rect(master)

    # Component 1: Non-placeholder rectangle shape exists on master slide (0.3 points)
    try:
        if len(decorative_shapes) > 0:
            print(f"PASS: Component 1 -- Found {len(decorative_shapes)} non-placeholder shape(s) on master slide (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 -- No non-placeholder shapes found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if len(decorative_shapes) == 0:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Find the best candidate: a rectangle-like shape near the bottom
    # Pick the shape closest to the bottom of the slide
    best_shape = None
    best_bottom = -1
    for shape in decorative_shapes:
        bottom = shape.top + shape.height
        if bottom > best_bottom:
            best_bottom = bottom
            best_shape = shape

    shape = best_shape
    print(f"  Candidate shape: name='{shape.name}', left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")

    # Component 2: Correct position and dimensions (0.4 points)
    # Expected: left=0, top~7.2in (6583680 EMU), width=10in (9144000 EMU), height~0.3in (274320 EMU)
    try:
        expected_width = slide_width  # full width = 10 inches
        expected_height = Inches(0.3)  # 0.3 inches = 274320 EMU
        expected_top = slide_height - expected_height  # bottom of slide minus bar height

        width_ok = is_approx_equal(shape.width, expected_width, tolerance=0.02)
        height_ok = is_approx_equal(shape.height, expected_height, tolerance=0.05)
        left_ok = shape.left <= Inches(0.1)  # should be at or very near left=0
        top_ok = is_approx_equal(shape.top, expected_top, tolerance=0.03)

        checks_passed = sum([width_ok, height_ok, left_ok, top_ok])
        print(f"  Position checks: width_ok={width_ok} (actual={shape.width}, expected={expected_width})")
        print(f"  height_ok={height_ok} (actual={shape.height}, expected={expected_height})")
        print(f"  left_ok={left_ok} (actual={shape.left})")
        print(f"  top_ok={top_ok} (actual={shape.top}, expected={expected_top})")

        if checks_passed == 4:
            print(f"PASS: Component 2 -- All position/dimension checks passed (0.4 pts)")
            total_score += 0.4
        elif checks_passed >= 3:
            print(f"PARTIAL: Component 2 -- {checks_passed}/4 position checks passed (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Only {checks_passed}/4 position checks passed")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct solid fill color #FF6B35 (0.3 points)
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID fill
            actual_rgb = str(fill.fore_color.rgb).upper()
            expected_rgb = "FF6B35"
            if actual_rgb == expected_rgb:
                print(f"PASS: Component 3 -- Fill color is #{actual_rgb} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 -- Fill color is #{actual_rgb}, expected #{expected_rgb}")
        else:
            fill_type = fill.type if fill.type is not None else "None"
            print(f"FAIL: Component 3 -- Fill type is {fill_type}, expected SOLID (1)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
