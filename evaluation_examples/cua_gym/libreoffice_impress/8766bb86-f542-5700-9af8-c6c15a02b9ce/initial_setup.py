"""
Initial Setup: Event Agenda presentation with 8 slides, no decorative bar on master.
Task ID: impress_ma_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard 10x7.5 slide dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Annual Tech Conference 2025"
    slide1.placeholders[1].text = "Event Agenda & Schedule\nSan Francisco Convention Center\nMarch 15-17, 2025"

    # --- Slide 2: Welcome & Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Welcome & Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Three days of innovation, networking, and learning"
    p = body2.add_paragraph()
    p.text = "Over 50 speakers from leading tech companies"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Workshops, keynotes, and panel discussions"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Networking events each evening"
    p.level = 1

    # --- Slide 3: Day 1 Schedule ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Day 1 - March 15"
    body3 = slide3.placeholders[1].text_frame
    schedule_day1 = [
        "9:00 AM - Registration & Coffee",
        "10:00 AM - Opening Keynote: Dr. Emily Watson",
        "11:30 AM - AI & Machine Learning Track",
        "12:30 PM - Lunch Break",
        "2:00 PM - Cloud Infrastructure Workshop",
        "4:00 PM - Lightning Talks",
        "6:00 PM - Welcome Reception",
    ]
    body3.text = schedule_day1[0]
    for item in schedule_day1[1:]:
        p = body3.add_paragraph()
        p.text = item

    # --- Slide 4: Day 2 Schedule ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Day 2 - March 16"
    body4 = slide4.placeholders[1].text_frame
    schedule_day2 = [
        "9:00 AM - Morning Networking Breakfast",
        "10:00 AM - Keynote: Sarah Chen, CTO of NovaTech",
        "11:30 AM - Security & Privacy Panel",
        "12:30 PM - Lunch & Expo Floor",
        "2:00 PM - DevOps Best Practices",
        "3:30 PM - Startup Pitch Competition",
        "6:30 PM - Gala Dinner",
    ]
    body4.text = schedule_day2[0]
    for item in schedule_day2[1:]:
        p = body4.add_paragraph()
        p.text = item

    # --- Slide 5: Day 3 Schedule ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Day 3 - March 17"
    body5 = slide5.placeholders[1].text_frame
    schedule_day3 = [
        "9:00 AM - Unconference Sessions",
        "10:30 AM - Future of Web Development",
        "12:00 PM - Closing Keynote: Marcus Johnson",
        "1:00 PM - Farewell Lunch",
        "2:30 PM - Hackathon Kickoff (Optional)",
    ]
    body5.text = schedule_day3[0]
    for item in schedule_day3[1:]:
        p = body5.add_paragraph()
        p.text = item

    # --- Slide 6: Keynote Speakers ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Featured Speakers"
    body6 = slide6.placeholders[1].text_frame
    speakers = [
        "Dr. Emily Watson - AI Research Lead, DeepMind",
        "Sarah Chen - CTO, NovaTech Systems",
        "Marcus Johnson - VP Engineering, Cloudflare",
        "Priya Patel - Director of Product, Stripe",
        "James O'Brien - Security Architect, CrowdStrike",
        "Lisa Nakamura - Head of DevRel, GitHub",
    ]
    body6.text = speakers[0]
    for s in speakers[1:]:
        p = body6.add_paragraph()
        p.text = s

    # --- Slide 7: Venue Information ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Venue & Logistics"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "San Francisco Convention Center, Hall A-C"
    items7 = [
        "Wi-Fi: TechConf2025 / Password: innovate2025",
        "Parking: Lot B with complimentary shuttle service",
        "Meals: Breakfast and lunch included with registration",
        "Emergency contact: +1 (415) 555-0192",
    ]
    for item in items7:
        p = body7.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You!"
    slide8.placeholders[1].text = "See you at TechConf 2026!\nwww.techconf2025.com\n#TechConf2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
