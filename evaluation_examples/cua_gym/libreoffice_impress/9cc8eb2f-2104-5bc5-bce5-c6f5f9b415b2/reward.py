"""
FINAL REWARD SCRIPT - SUCCESS
Task: Replace hyphens between number ranges with en dashes across the document.
Generated: 2025-10-17 11:20:41
Status: success
Model: azure-o3
Total Steps: 1
"""

import re
import os
from pptx import Presentation


def count_number_ranges(text: str):
    """Return counts of number-ranges written with hyphen vs. en-dash.

    A "number range" is any two numbers separated by either a hyphen (-) or an
    en dash (–) with optional surrounding spaces.
    """
    # hyphen between numbers (e.g. 1990-2000 or 1 - 5)
    hyphen_pattern = re.compile(r"(?<![\d])\d+[\d,.]*\s*-\s*\d+[\d,.]*(?![\d])")
    # en dash between numbers (–, U+2013)
    en_dash_pattern = re.compile(r"(?<![\d])\d+[\d,.]*\s*–\s*\d+[\d,.]*(?![\d])")
    hyphen_matches = hyphen_pattern.findall(text)
    en_dash_matches = en_dash_pattern.findall(text)
    return len(hyphen_matches), len(en_dash_matches)


def verify_replace_hyphen_with_en_dash(file_path: str) -> float:
    """Verify that all hyphens within number ranges have been replaced by en dashes.

    Scoring:
      • 1.0  – Every number-range uses an en dash and none use a hyphen.
      • <1.0 – Proportional to fraction of ranges correctly using en dash.
      • 0.0  – File missing / unreadable or no ranges found.
    """
    print(f"Checking presentation: {file_path}")

    # ---------- 0. Prerequisite checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Failed to load PPTX: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 1. Collect all text from every slide ----------
    all_text_parts = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text
                if txt:
                    slide_text_parts.append(txt)
        combined_slide_text = "\n".join(slide_text_parts)
        print(f"Slide {slide_idx}: collected {len(combined_slide_text)} characters of text")
        all_text_parts.append(combined_slide_text)

    document_text = "\n".join(all_text_parts)

    # ---------- 2. Count hyphen-ranges vs. en-dash-ranges ----------
    hyphen_cnt, en_dash_cnt = count_number_ranges(document_text)
    total_ranges = hyphen_cnt + en_dash_cnt
    print(f"Total number ranges found: {total_ranges}")
    print(f"  Hyphen ranges (-): {hyphen_cnt}")
    print(f"  En dash ranges (–): {en_dash_cnt}")

    # ---------- 3. Scoring ----------
    if total_ranges == 0:
        # Cannot verify task without any number range; give 0.
        print("✗ No number ranges detected; unable to verify replacement task.")
        print("REWARD: 0.0")
        return 0.0

    # Fraction of ranges correctly using en dash
    fraction_correct = en_dash_cnt / total_ranges
    print(f"Fraction of correct en dash usage: {fraction_correct:.2%}")

    # Full credit only when every range uses en dash (no hyphen ranges remain)
    final_score = 1.0 if hyphen_cnt == 0 else fraction_correct

    # Safety clamp
    final_score = max(0.0, min(1.0, final_score))

    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/replace_hyphens_between_number_ranges_with_en_dashes_across_the_document.pptx"
    verify_replace_hyphen_with_en_dash(FILE_PATH)

