"""
Initial Setup: Trade show booth display presentation with 8 slides, default settings.
Task ID: impress_gf2_050
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_050'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette for trade show
    DARK_BLUE = RGBColor(0x0D, 0x25, 0x3F)
    ACCENT_BLUE = RGBColor(0x1B, 0x6E, 0xB5)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    ORANGE = RGBColor(0xE8, 0x6C, 0x00)

    # ---- Slide 1: Title / Welcome ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_text_box(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
                 "NovaTech Solutions", font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(2), Inches(3.5), Inches(9), Inches(1),
                 "Innovating the Future of Industrial Automation",
                 font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(3), Inches(5.5), Inches(7), Inches(0.8),
                 "Visit us at Booth #4217  |  AutomationExpo 2025  |  Las Vegas Convention Center",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Company Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide2.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "About NovaTech Solutions", font_size=36, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.LEFT)
    overview_text = (
        "Founded in 2008, NovaTech Solutions has grown from a small engineering "
        "startup in Austin, TX to a global leader in industrial automation with "
        "over 2,400 employees across 18 countries.\n\n"
        "Our mission: Empower manufacturers with intelligent automation that "
        "reduces costs by 35% and increases throughput by 60%.\n\n"
        "Key Markets: Automotive, Aerospace, Semiconductor, Pharmaceutical"
    )
    add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                 overview_text, font_size=20, color=DARK_GRAY)

    # ---- Slide 3: Product Line - RoboAssist Pro ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide3.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "RoboAssist Pro 5000", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                 "Next-generation collaborative robot arm\n\n"
                 "• 6-axis articulation with 0.02mm precision\n"
                 "• Payload capacity: 25 kg\n"
                 "• Built-in vision system with AI object recognition\n"
                 "• Safety-rated force limiting (ISO 10218-1)\n"
                 "• Setup time: under 30 minutes\n\n"
                 "Starting at $47,500",
                 font_size=18, color=DARK_GRAY)
    add_text_box(slide3, Inches(7), Inches(2.5), Inches(5), Inches(1),
                 "[Product Image Placeholder]", font_size=14, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # ---- Slide 4: Product Line - SmartConveyor ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide4.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "SmartConveyor X-Series", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide4, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5),
                 "Intelligent material handling system\n\n"
                 "• Modular design — reconfigure in hours, not days\n"
                 "• IoT-enabled with real-time throughput analytics\n"
                 "• Energy-efficient servo motors (40% less power)\n"
                 "• Automatic jam detection and self-clearing\n"
                 "• Supports loads from 500g to 200kg\n\n"
                 "Starting at $32,000 per 10m section",
                 font_size=18, color=DARK_GRAY)
    add_text_box(slide4, Inches(7), Inches(2.5), Inches(5), Inches(1),
                 "[Product Image Placeholder]", font_size=14, color=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # ---- Slide 5: Case Study ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide5.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "Case Study: Meridian Automotive", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide5, Inches(0.8), Inches(1.8), Inches(11), Inches(4.5),
                 "Challenge: Meridian's stamping plant in Detroit was experiencing "
                 "18% downtime due to manual material handling bottlenecks.\n\n"
                 "Solution: Deployed 12 RoboAssist Pro units integrated with "
                 "SmartConveyor X-Series across 3 production lines.\n\n"
                 "Results after 6 months:\n"
                 "  • Downtime reduced from 18% to 3.2%\n"
                 "  • Throughput increased by 62%\n"
                 "  • ROI achieved in 14 months\n"
                 "  • Zero safety incidents since deployment\n\n"
                 "\"NovaTech transformed our operation.\" — James Kowalski, VP Manufacturing",
                 font_size=18, color=DARK_GRAY)

    # ---- Slide 6: Innovation Roadmap ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide6.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
                 "2025-2026 Innovation Roadmap", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(slide6, Inches(0.8), Inches(1.8), Inches(5), Inches(4),
                 "Q3 2025\n"
                 "  RoboAssist Pro 6000 — AI-driven path planning\n\n"
                 "Q4 2025\n"
                 "  SmartConveyor X2 — Predictive maintenance\n\n"
                 "Q1 2026\n"
                 "  NovaTech Cloud — Unified fleet management\n\n"
                 "Q2 2026\n"
                 "  AutoInspect — Computer vision quality control",
                 font_size=18, color=DARK_GRAY)
    add_text_box(slide6, Inches(6.5), Inches(1.8), Inches(6), Inches(4),
                 "Key R&D Investments:\n\n"
                 "  • $45M in AI / machine learning\n"
                 "  • $28M in advanced sensor tech\n"
                 "  • $18M in cloud infrastructure\n"
                 "  • Partnership with MIT Robotics Lab\n"
                 "  • Partnership with Fraunhofer IPA",
                 font_size=18, color=DARK_GRAY)

    # ---- Slide 7: Live Demo Schedule ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_text_box(slide7, Inches(1.5), Inches(0.5), Inches(10), Inches(1),
                 "Live Demonstrations — Booth #4217", font_size=36, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, Inches(1.5), Inches(2), Inches(10), Inches(4.5),
                 "10:00 AM — RoboAssist Pro Pick-and-Place Challenge\n\n"
                 "11:30 AM — SmartConveyor Live Reconfiguration\n\n"
                 "1:00 PM — AI Vision System Object Recognition\n\n"
                 "2:30 PM — Full Integration Demo: Robot + Conveyor + Vision\n\n"
                 "4:00 PM — Q&A with NovaTech Engineering Team",
                 font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ---- Slide 8: Contact / CTA ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_text_box(slide8, Inches(2), Inches(1), Inches(9), Inches(1.5),
                 "Let's Build the Future Together", font_size=42, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide8, Inches(2), Inches(3), Inches(9), Inches(3),
                 "Schedule a personalized demo at our booth\n\n"
                 "sales@novatechsolutions.com\n"
                 "+1 (512) 555-0198\n\n"
                 "www.novatechsolutions.com\n\n"
                 "Scan the QR code at our booth for a free automation assessment",
                 font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Save — NO transitions, NO auto-advance, NO kiosk mode
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
