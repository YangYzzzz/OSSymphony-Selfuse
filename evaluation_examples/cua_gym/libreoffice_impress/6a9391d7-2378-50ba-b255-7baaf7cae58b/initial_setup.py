"""
Initial Setup: Create a presentation with 8 slides, slide 3 has a column chart.
Task ID: impress_tct_070
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_070'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Chart data shared between initial and golden
CATEGORIES = ['Q1 2024', 'Q2 2024', 'Q3 2024', 'Q4 2024', 'Q1 2025']
SERIES1_NAME = 'North Region'
SERIES1_DATA = [42500, 38700, 51200, 47800, 55100]
SERIES2_NAME = 'South Region'
SERIES2_DATA = [31800, 44300, 36900, 52100, 48600]


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Regional Sales Performance Review"
    slide1.placeholders[1].text = "Fiscal Year 2024-2025 Analysis\nPrepared by Analytics Division"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Overview of Regional Performance"
    for item in ["Quarterly Revenue Breakdown", "Year-over-Year Trends",
                  "Key Growth Drivers", "Strategic Recommendations",
                  "Next Steps and Action Items", "Q&A Session"]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Column Chart (5 categories, 2 series) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Add title text box
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Quarterly Revenue by Region"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add column chart
    chart_data = CategoryChartData()
    chart_data.categories = CATEGORIES
    chart_data.add_series(SERIES1_NAME, SERIES1_DATA)
    chart_data.add_series(SERIES2_NAME, SERIES2_DATA)

    chart_shape = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(1.5), Inches(10), Inches(5.0),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True

    # --- Slide 4: Market Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Market Analysis"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North Region showed consistent growth across all quarters"
    for item in ["South Region experienced strong Q4 recovery (+41%)",
                  "Combined revenue exceeded targets by 12%",
                  "New client acquisitions drove Q3 North Region surge",
                  "Supply chain improvements reduced operational costs"]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Customer Insights ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Customer Insights"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Enterprise segment contributed 62% of total revenue"
    for item in ["SMB growth rate: 28% year-over-year",
                  "Customer retention improved to 94.2%",
                  "Net Promoter Score increased from 67 to 74",
                  "Average deal size grew by 15% in North Region"]:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 6: Operational Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Operational Metrics"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Average sales cycle reduced from 45 to 37 days"
    for item in ["Proposal-to-close ratio improved to 34%",
                  "Cross-sell revenue increased by $2.3M",
                  "Team headcount grew by 18% with maintained productivity",
                  "Training investment per rep: $4,200 annually"]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 7: Title Only (empty except title) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf7 = txBox7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Revenue Trends Visualization"
    p7.alignment = PP_ALIGN.LEFT
    run7 = p7.runs[0]
    run7.font.size = Pt(28)
    run7.font.bold = True
    run7.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps & Action Items"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Finalize Q1 2025 territory realignment by March 15"
    for item in ["Launch customer success pilot program in South Region",
                  "Submit budget proposals for FY2025-2026",
                  "Schedule quarterly business reviews with top 20 accounts",
                  "Deploy new CRM analytics dashboard by end of April"]:
        p = body8.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
