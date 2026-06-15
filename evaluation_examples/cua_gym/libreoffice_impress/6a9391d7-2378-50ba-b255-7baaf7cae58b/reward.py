"""
Reward Script: Duplicate chart from slide 3 to slide 7, change duplicated chart type from column to line
Task ID: impress_tct_070
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 7 contains a chart shape
  Component 2 (0.25): Slide 7 chart is LINE type (not column)
  Component 3 (0.25): Slide 7 chart data matches slide 3 data (same categories and series values)
  Component 4 (0.25): Slide 3 chart is still COLUMN_CLUSTERED (unchanged) AND slide 7 has a chart
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_070'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"PRECONDITION FAIL: Expected at least 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed
    slide7 = prs.slides[6]  # 0-indexed

    # Helper: find chart shape in a slide
    def find_chart_shape(slide):
        for shape in slide.shapes:
            if hasattr(shape, 'chart'):
                try:
                    _ = shape.chart.chart_type
                    return shape
                except Exception:
                    pass
        return None

    # Helper: extract chart data (categories + series values)
    def extract_chart_data(chart):
        try:
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            series_data = []
            for series in chart.series:
                series_data.append(list(series.values))
            return categories, series_data
        except Exception as e:
            print(f"  WARNING: Could not extract chart data: {e}")
            return None, None

    slide7_chart_shape = find_chart_shape(slide7)
    slide3_chart_shape = find_chart_shape(slide3)

    # Component 1: Slide 7 contains a chart shape (0.25 points)
    # This FAILS on initial (slide 7 has no chart) -> PASSES on golden
    try:
        if slide7_chart_shape is not None:
            print(f"PASS: Component 1 -- Slide 7 has a chart shape (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Slide 7 has no chart shape")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 7 chart type is LINE (0.25 points)
    # This FAILS on initial (no chart on slide 7) -> PASSES on golden (LINE chart)
    try:
        if slide7_chart_shape is not None:
            chart7 = slide7_chart_shape.chart
            chart7_type = chart7.chart_type
            # XL_CHART_TYPE.LINE = 4
            if chart7_type == XL_CHART_TYPE.LINE:
                print(f"PASS: Component 2 -- Slide 7 chart type is LINE ({chart7_type}) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 -- Slide 7 chart type is {chart7_type}, expected LINE (4)")
        else:
            print(f"FAIL: Component 2 -- No chart on slide 7 to check type")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 7 chart data matches slide 3 chart data (0.25 points)
    # Same categories (5) and same series values (2 series)
    # This FAILS on initial (no chart on slide 7) -> PASSES on golden (same data)
    try:
        if slide7_chart_shape is not None and slide3_chart_shape is not None:
            cats3, series3 = extract_chart_data(slide3_chart_shape.chart)
            cats7, series7 = extract_chart_data(slide7_chart_shape.chart)

            if cats3 is None or cats7 is None:
                print(f"FAIL: Component 3 -- Could not extract chart data for comparison")
            else:
                cats_match = cats3 == cats7
                series_match = len(series3) == len(series7)
                if series_match:
                    for i in range(len(series3)):
                        if series3[i] != series7[i]:
                            series_match = False
                            break

                if cats_match and series_match:
                    print(f"PASS: Component 3 -- Slide 7 chart data matches slide 3 data "
                          f"({len(cats7)} categories, {len(series7)} series) (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 3 -- Data mismatch: "
                          f"cats_match={cats_match}, series_match={series_match}")
                    if not cats_match:
                        print(f"  Slide 3 categories: {cats3}")
                        print(f"  Slide 7 categories: {cats7}")
                    if not series_match:
                        print(f"  Slide 3 series: {series3}")
                        print(f"  Slide 7 series: {series7}")
        else:
            print(f"FAIL: Component 3 -- Missing chart on slide 3 or slide 7")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 3 chart is still COLUMN_CLUSTERED AND slide 7 has a chart (0.25 points)
    # The compound condition ensures this FAILS on initial (slide 7 has no chart)
    # and PASSES on golden (slide 3 unchanged + slide 7 has chart)
    try:
        if slide3_chart_shape is not None and slide7_chart_shape is not None:
            chart3 = slide3_chart_shape.chart
            chart3_type = chart3.chart_type
            # XL_CHART_TYPE.COLUMN_CLUSTERED = 51
            if chart3_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                print(f"PASS: Component 4 -- Slide 3 chart is still COLUMN_CLUSTERED ({chart3_type}) "
                      f"and slide 7 has a chart (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 4 -- Slide 3 chart type changed to {chart3_type}, "
                      f"expected COLUMN_CLUSTERED (51)")
        else:
            if slide3_chart_shape is None:
                print(f"FAIL: Component 4 -- No chart on slide 3")
            else:
                print(f"FAIL: Component 4 -- No chart on slide 7 (compound condition)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
