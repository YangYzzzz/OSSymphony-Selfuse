"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 1, I’ve inserted an image named “Picture 1.” I want it to work as a full-bleed cover, so I need to enlarge it proportionally until it fills the whole slide and make sure it stays dead-center. How do I do that in LibreOffice Impress?
Generated: 2025-09-10 13:27:20
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def verify_full_bleed_image(file_path):
    """Verify that slide 1 contains an image called "Picture 1" (or any picture if
    the name is missing), enlarged proportionally so it fully covers the slide
    (full-bleed) and is centred on the canvas.

    Scoring (progressive – max 1.0):
        • 0.2 – required picture found on slide 1
        • 0.2 – width ≥ 98 % of slide width (≤ 105 % overspill allowed)
        • 0.2 – height ≥ 98 % of slide height (≤ 105 % overspill allowed)
        • 0.2 – horizontally centred (≤ 1 % deviation)
        • 0.2 – vertically centred   (≤ 1 % deviation)
    """

    print(f"Verifying presentation: {file_path}")
    score = 0.0
    max_score = 1.0

    # ---------- Basic file checks ----------
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        return 0.0

    if not prs.slides:
        print("✗ Presentation contains no slides.")
        return 0.0

    slide = prs.slides[0]  # First slide (slide 1)
    slide_w, slide_h = prs.slide_width, prs.slide_height  # in EMUs
    print(f"Slide size (EMU): width={slide_w}, height={slide_h}")

    # ---------- Locate the picture ----------
    picture = None
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and (shp.name or "").strip().lower() == "picture 1":
            picture = shp
            break
    if picture is None:
        # Fallback: take first picture on the slide
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture = shp
                break
    if picture is None:
        print("✗ No picture found on slide 1.")
        return 0.0

    print(
        f"✓ Picture found – name: '{picture.name}', left={picture.left}, top={picture.top}, "
        f"width={picture.width}, height={picture.height}")
    score += 0.2  # picture present

    # ---------- Coverage checks ----------
    # Must cover at least 98 % of slide dimension, but allow up to 105 % overspill.
    width_ok  = slide_w * 0.98 <= picture.width  <= slide_w * 1.05
    height_ok = slide_h * 0.98 <= picture.height <= slide_h * 1.05

    if width_ok:
        print("✓ Width covers slide adequately.")
        score += 0.2
    else:
        print("✗ Width does not sufficiently cover the slide.")

    if height_ok:
        print("✓ Height covers slide adequately.")
        score += 0.2
    else:
        print("✗ Height does not sufficiently cover the slide.")

    # ---------- Centring checks ----------
    horiz_diff = abs((picture.left + picture.width / 2) - (slide_w / 2))
    vert_diff  = abs((picture.top  + picture.height / 2) - (slide_h / 2))
    horiz_ok = horiz_diff <= slide_w * 0.01  # ≤ 1 % deviation
    vert_ok  = vert_diff  <= slide_h * 0.01

    if horiz_ok:
        print("✓ Image horizontally centred.")
        score += 0.2
    else:
        print(f"✗ Image not horizontally centred (diff={horiz_diff}).")

    if vert_ok:
        print("✓ Image vertically centred.")
        score += 0.2
    else:
        print(f"✗ Image not vertically centred (diff={vert_diff}).")

    # ---------- Final score ----------
    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------------------------------------------------------------------
# Execute verification (path provided by task environment) -------------------
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    PRESENTATION_PATH = "/home/user/on_slide_1_ive_inserted_an_image_named_picture_1_i_want_it_to_work_as_a_full_bleed_cover_so_i_need_t_golden.pptx"
    verify_full_bleed_image(PRESENTATION_PATH)
