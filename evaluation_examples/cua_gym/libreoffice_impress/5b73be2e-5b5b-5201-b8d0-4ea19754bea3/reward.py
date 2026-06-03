"""
Reward Script: Crop image on slide 4 to show only center portion (~15% from each edge)
Task ID: impress_fix_084
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): srcRect crop element exists on slide 4 image with non-zero crop values
  Component 2 (0.4): Each edge crop value is approximately 15% (within 10%-20% tolerance)
  Component 3 (0.3): Crop is roughly symmetric (all 4 edges within 5% of each other)
"""

import os
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_084'

def verify_task(file_path):
    """
    Verify that the image on slide 4 has been cropped with ~15% removed from each edge.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    # Find the picture shape on slide 4 (index 3)
    slide = prs.slides[3]
    pic_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_shape = shape
            break

    if pic_shape is None:
        print("FAIL: No picture shape found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Extract crop values from the XML srcRect element
    # srcRect values are in 1/1000ths of a percent (e.g., 15000 = 15%)
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    pic_elem = pic_shape._element
    blipFill = pic_elem.find(f'{{{ns_p}}}blipFill')
    if blipFill is None:
        blipFill = pic_elem.find(f'.//{{{ns_a}}}blipFill')

    srcRect = None
    if blipFill is not None:
        srcRect = blipFill.find(f'{{{ns_a}}}srcRect')

    crop_values = {}
    if srcRect is not None:
        for edge in ['l', 'r', 't', 'b']:
            val_str = srcRect.get(edge, '0')
            try:
                crop_values[edge] = int(val_str)
            except (ValueError, TypeError):
                crop_values[edge] = 0
    else:
        crop_values = {'l': 0, 'r': 0, 't': 0, 'b': 0}

    print(f"INFO: Crop values (1/1000ths of %): l={crop_values['l']}, r={crop_values['r']}, t={crop_values['t']}, b={crop_values['b']}")

    # Component 1: srcRect exists with non-zero crop values on all 4 edges (0.3 points)
    # This checks that cropping was actually applied (changes from initial state which has no crop)
    try:
        all_nonzero = all(crop_values[e] > 0 for e in ['l', 'r', 't', 'b'])
        if srcRect is not None and all_nonzero:
            print(f"PASS: Component 1 — srcRect exists with non-zero crop on all 4 edges (0.3 pts)")
            total_score += 0.3
        else:
            if srcRect is None:
                print(f"FAIL: Component 1 — No srcRect element found (no crop applied)")
            else:
                zero_edges = [e for e in ['l', 'r', 't', 'b'] if crop_values[e] == 0]
                print(f"FAIL: Component 1 — Zero crop on edges: {zero_edges}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Each edge crop is approximately 15% (within 10%-20% range) (0.4 points)
    # 15% = 15000 in 1/1000ths of percent. Accept 10000-20000 (10%-20%)
    try:
        edges_ok = 0
        for edge_name, edge_key in [('left', 'l'), ('right', 'r'), ('top', 't'), ('bottom', 'b')]:
            val = crop_values[edge_key]
            pct = val / 1000.0  # convert to percentage
            if 10.0 <= pct <= 20.0:
                edges_ok += 1
                print(f"  {edge_name}: {pct:.1f}% — within acceptable range")
            else:
                print(f"  {edge_name}: {pct:.1f}% — outside acceptable range [10%-20%]")

        if edges_ok == 4:
            print(f"PASS: Component 2 — All 4 edges cropped approximately 15% (0.4 pts)")
            total_score += 0.4
        elif edges_ok >= 2:
            partial = 0.4 * (edges_ok / 4)
            print(f"PARTIAL: Component 2 — {edges_ok}/4 edges in range ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Only {edges_ok}/4 edges in acceptable range")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Crop is roughly symmetric — all 4 edges within 5 percentage points of each other (0.3 points)
    # This ensures the "center portion" requirement is met (not just arbitrary crop)
    try:
        vals_pct = [crop_values[e] / 1000.0 for e in ['l', 'r', 't', 'b']]
        if all(v > 0 for v in vals_pct):
            max_diff = max(vals_pct) - min(vals_pct)
            if max_diff <= 5.0:
                print(f"PASS: Component 3 — Crop is symmetric (max difference: {max_diff:.1f}%) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Crop not symmetric (max difference: {max_diff:.1f}%, need <= 5%)")
        else:
            print(f"FAIL: Component 3 — Not all edges have positive crop")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
