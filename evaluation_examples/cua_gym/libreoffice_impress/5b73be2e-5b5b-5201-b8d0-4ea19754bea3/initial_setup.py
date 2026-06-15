"""
Initial Setup: Create a presentation with a large image on slide 4 (no cropping).
Task ID: impress_fix_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import struct
import zlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/_scene_photo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_gradient_png(path, width=960, height=720):
    """Create a colorful gradient PNG image (a scenic-looking pattern)."""
    # Build raw RGBA pixel data with a multi-color gradient pattern
    rows = []
    for y in range(height):
        row = b''
        for x in range(width):
            r = int(40 + 180 * (x / width))
            g = int(60 + 160 * (y / height))
            b = int(180 - 120 * (x / width) + 40 * (y / height))
            b = max(0, min(255, b))
            row += struct.pack('BBB', r, g, b)
        rows.append(row)

    # Manually build a PNG file
    def make_chunk(chunk_type, data):
        chunk = chunk_type + data
        return struct.pack('>I', len(data)) + chunk + struct.pack('>I', zlib.crc32(chunk) & 0xFFFFFFFF)

    signature = b'\x89PNG\r\n\x1a\n'
    ihdr_data = struct.pack('>IIBBBBB', width, height, 8, 2, 0, 0, 0)  # 8-bit RGB
    ihdr = make_chunk(b'IHDR', ihdr_data)

    raw_data = b''
    for row in rows:
        raw_data += b'\x00' + row  # filter byte 0 (None) per row
    compressed = zlib.compress(raw_data, 6)
    idat = make_chunk(b'IDAT', compressed)
    iend = make_chunk(b'IEND', b'')

    with open(path, 'wb') as f:
        f.write(signature + ihdr + idat + iend)


def create_initial():
    # Create the image file first
    create_gradient_png(IMG_PATH, 960, 720)

    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Focus Image Workshop"
    slide1.placeholders[1].text = "Techniques for Visual Composition"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Introduction to visual focus techniques"
    items = [
        "Understanding the rule of thirds",
        "Cropping for emphasis and storytelling",
        "Practical exercises with sample images",
        "Review and feedback session",
    ]
    for item in items:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 3: Key Concepts ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Concepts in Visual Focus"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Drawing the viewer's eye to the subject"
    concepts = [
        "Eliminate distracting elements at the edges",
        "Use cropping to tighten composition",
        "Balance negative space with the focal point",
        "Consider aspect ratio changes carefully",
    ]
    for c in concepts:
        p = tf3.add_paragraph()
        p.text = c
        p.level = 1

    # --- Slide 4: The full scene image (8in x 6in, no cropping) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Center the 8x6 inch image on the slide
    img_w = Inches(8)
    img_h = Inches(6)
    left = (prs.slide_width - img_w) // 2
    top = (prs.slide_height - img_h) // 2
    pic = slide4.shapes.add_picture(IMG_PATH, left, top, img_w, img_h)

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Summary & Next Steps"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Practice cropping with your own photos"
    steps = [
        "Select images with clear subjects",
        "Experiment with removing 10-20% from edges",
        "Compare before and after compositions",
        "Share results in the next session",
    ]
    for s in steps:
        p = tf5.add_paragraph()
        p.text = s
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp image
    os.remove(IMG_PATH)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
