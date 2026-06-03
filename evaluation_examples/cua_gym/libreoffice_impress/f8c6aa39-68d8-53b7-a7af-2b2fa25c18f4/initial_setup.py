"""
Initial Setup: 6-slide executive summary deck — all text regular weight, title sizes vary, no underlines
Task ID: osworld_impress_bold_all_title_size_underline_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_run(tf, text, font_size_pt, bold=False, italic=False, underline=False, color=None):
    """Helper: clear text frame and set first paragraph with a run."""
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if color:
        run.font.color.rgb = color


def add_body_paragraph(tf, text, font_size_pt=18, bold=False, italic=False):
    """Add a paragraph with a run to an existing text frame."""
    para = tf.add_paragraph()
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    return para


def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # -----------------------------------------------------------------------
    # Slide 1: Title slide — "Q2 2025 Executive Summary"
    # -----------------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout

    # Title placeholder
    title_ph = slide1.shapes.title
    title_ph.text = ""
    tf = title_ph.text_frame
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = "Q2 2025 Executive Summary"
    run.font.size = Pt(32)   # Varied — NOT 38pt
    run.font.bold = False
    run.font.underline = False  # Explicitly not underlined

    # Subtitle placeholder
    subtitle_ph = slide1.placeholders[1]
    subtitle_ph.text = ""
    tf_sub = subtitle_ph.text_frame
    para_sub = tf_sub.paragraphs[0]
    run_sub = para_sub.add_run()
    run_sub.text = "Prepared by: Strategy & Operations"
    run_sub.font.size = Pt(20)
    run_sub.font.bold = False

    para_sub2 = tf_sub.add_paragraph()
    run_sub2 = para_sub2.add_run()
    run_sub2.text = "June 30, 2025"
    run_sub2.font.size = Pt(18)
    run_sub2.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 2: Company Overview — title size 28pt
    # -----------------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

    tf2_title = slide2.shapes.title.text_frame
    tf2_title.paragraphs[0].text = ""
    p2t = tf2_title.paragraphs[0]
    r2t = p2t.add_run()
    r2t.text = "Company Overview"
    r2t.font.size = Pt(28)   # Varied — NOT 38pt
    r2t.font.bold = False

    tf2_body = slide2.placeholders[1].text_frame
    tf2_body.paragraphs[0].text = ""
    for bullet_text in [
        "Founded 2012 | Headquartered in San Francisco, CA",
        "1,240 full-time employees across 8 global offices",
        "Core verticals: SaaS, FinTech, Healthcare IT",
        "Revenue model: Annual Recurring Revenue (ARR) based",
        "Latest funding: Series D — $120M closed March 2025",
    ]:
        p = tf2_body.add_paragraph()
        r = p.add_run()
        r.text = bullet_text
        r.font.size = Pt(18)
        r.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 3: Financial Highlights — title size 30pt
    # -----------------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])

    tf3_title = slide3.shapes.title.text_frame
    tf3_title.paragraphs[0].text = ""
    p3t = tf3_title.paragraphs[0]
    r3t = p3t.add_run()
    r3t.text = "Financial Highlights"
    r3t.font.size = Pt(30)   # Varied — NOT 38pt
    r3t.font.bold = False

    tf3_body = slide3.placeholders[1].text_frame
    tf3_body.paragraphs[0].text = ""
    for bullet_text in [
        "ARR: $58.4M (+34% YoY)",
        "Gross Margin: 72% (up from 68% in Q2 2024)",
        "Net Revenue Retention: 118%",
        "New Bookings: $12.1M (Q2 2025)",
        "Churn Rate: 4.2% (all-time low)",
        "Operating Cash Flow: $3.8M positive",
    ]:
        p = tf3_body.add_paragraph()
        r = p.add_run()
        r.text = bullet_text
        r.font.size = Pt(18)
        r.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 4: Product Roadmap — title size 26pt
    # -----------------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])

    tf4_title = slide4.shapes.title.text_frame
    tf4_title.paragraphs[0].text = ""
    p4t = tf4_title.paragraphs[0]
    r4t = p4t.add_run()
    r4t.text = "Product Roadmap — H2 2025"
    r4t.font.size = Pt(26)   # Varied — NOT 38pt
    r4t.font.bold = False

    tf4_body = slide4.placeholders[1].text_frame
    tf4_body.paragraphs[0].text = ""
    for bullet_text in [
        "Q3: Launch AI-assisted reporting module (50 beta customers)",
        "Q3: Mobile app v3.0 — offline-first architecture",
        "Q4: Enterprise SSO & SCIM provisioning",
        "Q4: Expanded API marketplace (30+ third-party integrations)",
        "Q4: ISO 27001 certification audit",
    ]:
        p = tf4_body.add_paragraph()
        r = p.add_run()
        r.text = bullet_text
        r.font.size = Pt(18)
        r.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 5: People & Culture — title size 34pt
    # -----------------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])

    tf5_title = slide5.shapes.title.text_frame
    tf5_title.paragraphs[0].text = ""
    p5t = tf5_title.paragraphs[0]
    r5t = p5t.add_run()
    r5t.text = "People & Culture"
    r5t.font.size = Pt(34)   # Varied — NOT 38pt
    r5t.font.bold = False

    tf5_body = slide5.placeholders[1].text_frame
    tf5_body.paragraphs[0].text = ""
    for bullet_text in [
        "Headcount: +210 hires in H1 2025 (net +148)",
        "Engineering: 42% of total workforce",
        "Employee NPS: 72 (up 8 points from H2 2024)",
        "D&I: 38% women in leadership roles",
        "Remote-first policy: 65% distributed globally",
        "Annual learning budget: $3,000 per employee",
    ]:
        p = tf5_body.add_paragraph()
        r = p.add_run()
        r.text = bullet_text
        r.font.size = Pt(18)
        r.font.bold = False

    # -----------------------------------------------------------------------
    # Slide 6: Strategic Priorities — title size 29pt
    # -----------------------------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])

    tf6_title = slide6.shapes.title.text_frame
    tf6_title.paragraphs[0].text = ""
    p6t = tf6_title.paragraphs[0]
    r6t = p6t.add_run()
    r6t.text = "Strategic Priorities for H2 2025"
    r6t.font.size = Pt(29)   # Varied — NOT 38pt
    r6t.font.bold = False

    tf6_body = slide6.placeholders[1].text_frame
    tf6_body.paragraphs[0].text = ""
    for bullet_text in [
        "1. Accelerate enterprise segment — target 25 new Fortune 500 accounts",
        "2. International expansion — APAC go-to-market launch (Singapore HQ)",
        "3. AI-native product evolution — LLM integration across all modules",
        "4. Operational efficiency — reduce CAC by 15% via PLG motion",
        "5. Strategic M&A — evaluate 2-3 bolt-on acquisitions in BI/analytics",
    ]:
        p = tf6_body.add_paragraph()
        r = p.add_run()
        r.text = bullet_text
        r.font.size = Pt(18)
        r.font.bold = False

    # Save the file
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
