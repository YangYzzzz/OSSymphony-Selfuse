"""
Reward Script: Apply bold to all text, set all title font sizes to 38pt,
               and underline the title on slide 1 with a solid underline.
Task ID: osworld_impress_bold_all_title_size_underline_010
Domain: libreoffice_impress
Scoring:
  Component 1: All text runs in all slides are bold                    — 0.4 pts
  Component 2: All title shapes across all slides have 38pt font size  — 0.3 pts
  Component 3: Slide 1 title is underlined (solid underline)           — 0.3 pts
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_010'

EXPECTED_TITLE_SIZE_PT = 38.0
EXPECTED_TITLE_SIZE_EMU = int(EXPECTED_TITLE_SIZE_PT * 12700)  # 482600 EMU


def get_all_text_shapes(slide):
    """Recursively get all shapes with text frames, including nested group shapes."""
    def extract(shape):
        results = []
        if hasattr(shape, 'text_frame'):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def is_title_shape(shape):
    """
    Detect title shapes via placeholder type only.
    PP_PLACEHOLDER.TITLE = 1 (TITLE), PP_PLACEHOLDER.CENTER_TITLE = 3 (CENTER_TITLE).
    Explicitly excludes SUBTITLE (4) which is NOT a title placeholder.
    """
    try:
        ph_type = shape.placeholder_format.type
        # Only TITLE (1) and CENTER_TITLE (3) are "title" placeholders
        # SUBTITLE (4), OBJECT (7), BODY (2), etc. are excluded
        if str(ph_type) in ('TITLE (1)', 'CENTER_TITLE (3)'):
            return True
        # Fallback: check numeric value if enum repr differs
        ph_type_val = getattr(ph_type, 'real', None)
        if ph_type_val in (1, 3):
            return True
        # String fallback
        ph_str = str(ph_type)
        if ph_str.startswith('TITLE (') or ph_str.startswith('CENTER_TITLE ('):
            return True
    except Exception:
        pass
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides == 0:
        print("CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    # ----------------------------------------------------------------
    # Component 1: All text runs in all slides must be bold (0.4 points)
    # In initial_env: all runs have bold=False → FAIL → earns 0.0
    # In golden_env: all runs have bold=True   → PASS → earns 0.4
    # ----------------------------------------------------------------
    try:
        non_bold_count = 0
        total_run_count = 0
        non_bold_details = []

        for slide_idx, slide in enumerate(prs.slides):
            text_shapes = get_all_text_shapes(slide)
            for shape in text_shapes:
                for para in shape.text_frame.paragraphs:
                    nonempty = [r for r in para.runs if (r.text or "").strip()]
                    for run in nonempty:
                        total_run_count += 1
                        # Treat None as False (not bold)
                        bold = run.font.bold if run.font.bold is not None else False
                        if not bold:
                            non_bold_count += 1
                            if len(non_bold_details) < 5:
                                non_bold_details.append(
                                    f"Slide {slide_idx+1}, shape '{shape.name}', "
                                    f"run '{run.text[:20]}'"
                                )

        if total_run_count == 0:
            print("WARN: No non-empty text runs found in presentation")
        elif non_bold_count == 0:
            print(f"PASS: Component 1 — All {total_run_count} text runs are bold (0.4 pts)")
            total_score += 0.4
        else:
            print(
                f"FAIL: Component 1 — {non_bold_count}/{total_run_count} runs are NOT bold. "
                f"First issues: {non_bold_details}"
            )
    except Exception as e:
        print(f"ERROR: Component 1 (bold check) — {e}")

    # ----------------------------------------------------------------
    # Component 2: All title shapes across all slides have 38pt font size (0.3 points)
    # In initial_env: title sizes vary (not 38pt) → FAIL → earns 0.0
    # In golden_env: all titles are 38pt          → PASS → earns 0.3
    # ----------------------------------------------------------------
    try:
        wrong_size_count = 0
        title_count = 0
        wrong_size_details = []

        for slide_idx, slide in enumerate(prs.slides):
            text_shapes = get_all_text_shapes(slide)
            for shape in text_shapes:
                if is_title_shape(shape):
                    for para in shape.text_frame.paragraphs:
                        nonempty = [r for r in para.runs if (r.text or "").strip()]
                        for run in nonempty:
                            title_count += 1
                            run_size = run.font.size  # in EMU, or None
                            # Size None means inherited; check actual EMU value
                            if run_size is None or abs(run_size - EXPECTED_TITLE_SIZE_EMU) > 0:
                                wrong_size_count += 1
                                size_pt = (run_size / 12700) if run_size else None
                                if len(wrong_size_details) < 5:
                                    wrong_size_details.append(
                                        f"Slide {slide_idx+1}, shape '{shape.name}', "
                                        f"run '{run.text[:20]}', size={size_pt}pt (expected 38pt)"
                                    )

        if title_count == 0:
            print("WARN: No title shape runs found in presentation")
        else:
            comp2_pass = (wrong_size_count == 0)
            if comp2_pass:
                print(f"PASS: Component 2 — All {title_count} title runs have 38pt font size (0.3 pts)")
                total_score += 0.3
            else:
                print(
                    f"FAIL: Component 2 — {wrong_size_count}/{title_count} title runs do NOT have "
                    f"38pt. First issues: {wrong_size_details}"
                )
    except Exception as e:
        print(f"ERROR: Component 2 (title size check) — {e}")

    # ----------------------------------------------------------------
    # Component 3: Slide 1 title is underlined (solid underline) (0.3 points)
    # In initial_env: slide 1 title underline=False → FAIL → earns 0.0
    # In golden_env: slide 1 title underline=True   → PASS → earns 0.3
    # ----------------------------------------------------------------
    try:
        slide1 = prs.slides[0]
        title_shape = None
        for shape in get_all_text_shapes(slide1):
            if is_title_shape(shape):
                title_shape = shape
                break

        if title_shape is None:
            print("FAIL: Component 3 — No title shape found on slide 1")
        else:
            title_runs = []
            for para in title_shape.text_frame.paragraphs:
                nonempty = [r for r in para.runs if (r.text or "").strip()]
                title_runs.extend(nonempty)

            if not title_runs:
                print("FAIL: Component 3 — Title shape on slide 1 has no non-empty runs")
            else:
                non_underlined_count = 0
                for run in title_runs:
                    # underline=True means underlined; False or None means not underlined
                    underline = run.font.underline if run.font.underline is not None else False
                    if not underline:
                        non_underlined_count += 1

                comp3_pass = (non_underlined_count == 0)
                if comp3_pass:
                    total_score += 0.3
                    print(
                        f"PASS: Component 3 — Slide 1 title is underlined "
                        f"({len(title_runs)} run(s) checked) (0.3 pts)"
                    )
                else:
                    print(
                        f"FAIL: Component 3 — {non_underlined_count}/{len(title_runs)} "
                        f"slide 1 title runs are NOT underlined"
                    )
    except Exception as e:
        print(f"ERROR: Component 3 (underline check) — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
