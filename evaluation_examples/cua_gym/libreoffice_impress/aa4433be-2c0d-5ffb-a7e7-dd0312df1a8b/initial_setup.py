"""
Initial Setup: Apply strikethrough to lines 2 and 4 in the bullet list on slide 5.
Task ID: osworld_impress_strikethrough_text_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----- Slide 1: Title Slide -----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "TechNova Product Roadmap 2025"
    slide1.placeholders[1].text = "Strategic Initiatives & Milestones\nQ1 2025 – Q4 2025"

    # ----- Slide 2: Executive Summary -----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Drive 30% revenue growth through platform expansion"
    p = tf2.add_paragraph()
    p.text = "Launch three major product lines in key markets"
    p = tf2.add_paragraph()
    p.text = "Reduce customer churn by 15% with retention programs"
    p = tf2.add_paragraph()
    p.text = "Expand enterprise partnerships to 50+ accounts"

    # ----- Slide 3: Q1 Initiatives -----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Q1 2025 Initiatives"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Cloud Infrastructure Upgrade — Migrate to multi-region AWS"
    p = tf3.add_paragraph()
    p.text = "API Gateway Launch — Enable third-party integrations"
    p = tf3.add_paragraph()
    p.text = "Security Audit — SOC 2 Type II certification target"
    p = tf3.add_paragraph()
    p.text = "Mobile App Redesign — Improve onboarding UX flow"

    # ----- Slide 4: Q2 Milestones -----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Q2 2025 Milestones"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Analytics Dashboard v2.0 — Real-time reporting"
    p = tf4.add_paragraph()
    p.text = "Partner Portal Launch — Self-service onboarding"
    p = tf4.add_paragraph()
    p.text = "AI Feature Set — Predictive recommendations engine"
    p = tf4.add_paragraph()
    p.text = "Global CDN Rollout — Sub-50ms latency target"

    # ----- Slide 5: Feature Backlog & Prioritization -----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Feature Backlog & Prioritization"
    tf5 = slide5.placeholders[1].text_frame

    # 5 bullet items — NO strikethrough on any (initial state)
    bullets = [
        "Single Sign-On (SSO) integration across all products",
        "Legacy REST API deprecation and GraphQL migration",
        "Real-time collaboration editing in document workspace",
        "On-premise deployment option for enterprise clients",
        "Advanced role-based access control (RBAC) system",
    ]

    tf5.text = bullets[0]
    run0 = tf5.paragraphs[0].runs[0]
    # Ensure no strikethrough (explicitly set noStrike)
    run0.font._element.attrib['strike'] = 'noStrike'

    for bullet_text in bullets[1:]:
        p = tf5.add_paragraph()
        p.text = bullet_text
        run = p.runs[0]
        # Ensure no strikethrough
        run.font._element.attrib['strike'] = 'noStrike'

    # ----- Slide 6: Q3–Q4 Roadmap -----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Q3–Q4 2025 Roadmap"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Marketplace Launch — Third-party app ecosystem"
    p = tf6.add_paragraph()
    p.text = "Machine Learning Pipeline — AutoML model training"
    p = tf6.add_paragraph()
    p.text = "International Expansion — EMEA and APAC markets"
    p = tf6.add_paragraph()
    p.text = "Enterprise SLA Tier — 99.99% uptime guarantee"

    # ----- Slide 7: Key Metrics & Success Criteria -----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Metrics & Success Criteria"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Monthly Active Users: Target 500,000 by Q4 2025"
    p = tf7.add_paragraph()
    p.text = "Net Promoter Score (NPS): Achieve score of 65+"
    p = tf7.add_paragraph()
    p.text = "Revenue Run Rate: $18M ARR by end of fiscal year"
    p = tf7.add_paragraph()
    p.text = "Customer Retention Rate: Maintain 92% or above"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
