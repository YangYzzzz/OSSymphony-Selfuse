"""
Reward Script: Apply strikethrough to lines 2 and 4 in the bullet list on slide 5.
Task ID: osworld_impress_strikethrough_text_005
Domain: libreoffice_impress
Scoring:
  - Component 1: Line 2 (para index 1) on slide 5 has sngStrike applied       (0.4 pts)
  - Component 2: Line 4 (para index 3) on slide 5 has sngStrike applied       (0.4 pts)
  - Component 3: Lines 1, 3, 5 (para indices 0, 2, 4) do NOT have strikethrough (0.2 pts)
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_005'


def get_strike(run):
    """Return the strike attribute value for a run ('noStrike', 'sngStrike', 'dblStrike')."""
    return run.font._element.attrib.get('strike', 'noStrike')


def get_bullet_paras(slide):
    """
    Return the paragraphs of the content placeholder (bullet list) on slide 5.
    The content placeholder is 'Content Placeholder 2' (shape index 1).
    """
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name != 'Title 1':
            return shape.text_frame.paragraphs
    return []


def verify_task(file_path):
    """
    Verify that strikethrough has been applied to lines 2 and 4 of the
    bullet list on slide 5, and NOT applied to lines 1, 3, 5.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load file — precondition gate
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: check slide count is 7
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed: slide 5 is index 4

    # Retrieve bullet paragraphs from slide 5 content placeholder
    paras = get_bullet_paras(slide5)
    if len(paras) < 5:
        print(f"CRITICAL: Expected at least 5 bullet paragraphs on slide 5, found {len(paras)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Line 2 (para index 1) has strikethrough applied (0.4 points)
    try:
        line2_para = paras[1]
        line2_runs = [r for r in line2_para.runs if (r.text or "").strip()]
        if not line2_runs:
            print(f"FAIL: Component 1 — no non-empty runs found in line 2 (para index 1)")
        else:
            # All non-empty runs in the paragraph should have sngStrike
            line2_strikes = [get_strike(r) for r in line2_runs]
            if all(s in ('sngStrike', 'dblStrike') for s in line2_strikes):
                print(f"PASS: Component 1 — Line 2 has strikethrough: {line2_strikes} (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 1 — Line 2 expected strikethrough, found: {line2_strikes}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Line 4 (para index 3) has strikethrough applied (0.4 points)
    try:
        line4_para = paras[3]
        line4_runs = [r for r in line4_para.runs if (r.text or "").strip()]
        if not line4_runs:
            print(f"FAIL: Component 2 — no non-empty runs found in line 4 (para index 3)")
        else:
            # All non-empty runs in the paragraph should have sngStrike
            line4_strikes = [get_strike(r) for r in line4_runs]
            if all(s in ('sngStrike', 'dblStrike') for s in line4_strikes):
                print(f"PASS: Component 2 — Line 4 has strikethrough: {line4_strikes} (0.4 pts)")
                total_score += 0.4
            else:
                print(f"FAIL: Component 2 — Line 4 expected strikethrough, found: {line4_strikes}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Lines 2 AND 4 have strikethrough AND lines 1, 3, 5 do NOT (0.2 points)
    # This is a compound check: strikethrough applied precisely to exactly lines 2 and 4.
    # It FAILS on initial (where lines 2 and 4 have no strikethrough) and PASSES on golden.
    try:
        # Sub-check A: lines 2 and 4 must have strikethrough (already verified in components 1 and 2)
        line2_runs_all = [r for r in paras[1].runs if (r.text or "").strip()]
        line4_runs_all = [r for r in paras[3].runs if (r.text or "").strip()]
        line2_has_strike = line2_runs_all and all(get_strike(r) in ('sngStrike', 'dblStrike') for r in line2_runs_all)
        line4_has_strike = line4_runs_all and all(get_strike(r) in ('sngStrike', 'dblStrike') for r in line4_runs_all)

        if not (line2_has_strike and line4_has_strike):
            # Lines 2 and 4 don't have strikethrough; component 3 can't pass
            print(f"FAIL: Component 3 — prerequisite not met: lines 2 and 4 must have strikethrough")
        else:
            # Sub-check B: lines 1, 3, 5 must NOT have strikethrough
            unchanged_indices = [0, 2, 4]
            all_unchanged_ok = True
            for idx in unchanged_indices:
                para = paras[idx]
                runs = [r for r in para.runs if (r.text or "").strip()]
                if not runs:
                    continue
                for r in runs:
                    s = get_strike(r)
                    if s in ('sngStrike', 'dblStrike'):
                        print(f"FAIL: Component 3 — Line {idx+1} (para index {idx}) unexpectedly has strikethrough: {s}")
                        all_unchanged_ok = False
                        break
                if not all_unchanged_ok:
                    break

            if all_unchanged_ok:
                print(f"PASS: Component 3 — Strikethrough applied precisely (lines 2, 4 only; lines 1, 3, 5 unchanged) (0.2 pts)")
                total_score += 0.2
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
