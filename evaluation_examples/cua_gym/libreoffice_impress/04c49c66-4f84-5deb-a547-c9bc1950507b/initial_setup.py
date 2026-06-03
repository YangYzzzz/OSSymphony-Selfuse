"""
Initial Setup: Create a 9-slide Market Pitch presentation with slide 6 having empty content area
Task ID: impress_sales_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_textbox(slide, text, left, top, width, height,
                      font_size=28, bold=True, color=RGBColor(0x1A, 0x1A, 0x2E),
                      alignment=PP_ALIGN.LEFT):
    """Add a styled title text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_body_textbox(slide, text, left, top, width, height,
                     font_size=14, color=RGBColor(0x33, 0x33, 0x33)):
    """Add a body text box with the given content."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=14, color=RGBColor(0x33, 0x33, 0x33)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
    ACCENT_BLUE = RGBColor(0x00, 0x72, 0xC6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
    DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    add_title_textbox(slide1, "Market Pitch 2024", Inches(1), Inches(2), Inches(11), Inches(1.5),
                      font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_body_textbox(slide1, "Accelerating Growth Through Innovation",
                     Inches(1), Inches(3.8), Inches(11), Inches(1),
                     font_size=22, color=RGBColor(0xAA, 0xAA, 0xCC))
    add_body_textbox(slide1, "Presented by: Elena Rodriguez, VP of Sales  |  Q4 Strategy Review",
                     Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                     font_size=14, color=RGBColor(0x88, 0x88, 0xAA))

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide2, "Executive Summary", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    add_bullet_list(slide2, [
        "Revenue grew 18% YoY to $142M in FY2024, exceeding projections by $7M",
        "Customer acquisition cost reduced by 12% through optimized digital channels",
        "Net Promoter Score improved from 62 to 71, reflecting stronger customer loyalty",
        "Successfully launched 3 new product lines generating $28M in incremental revenue",
        "Expanded into 4 new international markets: Germany, Japan, Brazil, and Australia",
        "Strategic partnership with Meridian Corp. projected to unlock $15M pipeline in 2025",
    ], Inches(0.8), Inches(1.5), Inches(11), Inches(5), font_size=16, color=DARK_TEXT)

    # ---- Slide 3: Product Overview ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide3, "Product Overview", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    add_bullet_list(slide3, [
        "ProSuite Enterprise — flagship SaaS platform, 45% of revenue ($63.9M)",
        "DataSync Pro — real-time analytics module, adopted by 320+ enterprise clients",
        "SecureVault — compliance & data security layer, SOC 2 Type II certified",
        "MobileFirst SDK — cross-platform toolkit powering 1,200+ partner apps",
        "CloudBridge API — integration hub connecting 85+ third-party services",
    ], Inches(0.8), Inches(1.5), Inches(11), Inches(3), font_size=16, color=DARK_TEXT)
    add_body_textbox(slide3, "Product roadmap for 2025 includes AI-powered automation features and expanded API marketplace.",
                     Inches(0.8), Inches(5.0), Inches(11), Inches(1), font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ---- Slide 4: Revenue Highlights ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide4, "Revenue Highlights", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)

    # Revenue table
    table_shape = slide4.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(10), Inches(3.5))
    table = table_shape.table
    headers = ["Quarter", "Revenue ($M)", "Growth (%)", "New Customers"]
    data_rows = [
        ["Q1 2024", "$32.4", "15.2%", "187"],
        ["Q2 2024", "$34.8", "17.1%", "214"],
        ["Q3 2024", "$36.1", "18.9%", "239"],
        ["Q4 2024", "$38.7", "20.3%", "261"],
        ["FY2024 Total", "$142.0", "18.0%", "901"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '0072C6'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # ---- Slide 5: Customer Testimonials ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide5, "Customer Testimonials", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    testimonials = [
        ('"ProSuite transformed our workflow efficiency by 40%. Implementation was seamless and the ROI exceeded our expectations within 6 months."',
         "— Sarah Chen, CTO, NovaTech Industries"),
        ('"The DataSync Pro module gives us real-time visibility across all departments. We can now make data-driven decisions in minutes, not days."',
         "— Marcus Williams, VP Operations, Apex Financial Group"),
        ('"SecureVault gave us peace of mind during our SOC 2 audit. Compliance has never been easier to maintain at scale."',
         "— Priya Sharma, CISO, HealthNet Solutions"),
    ]
    y_pos = 1.5
    for quote, attribution in testimonials:
        add_body_textbox(slide5, quote, Inches(1.2), Inches(y_pos), Inches(10), Inches(0.9),
                         font_size=15, color=DARK_TEXT)
        add_body_textbox(slide5, attribution, Inches(1.5), Inches(y_pos + 0.9), Inches(10), Inches(0.4),
                         font_size=12, color=RGBColor(0x00, 0x72, 0xC6))
        y_pos += 1.7

    # ---- Slide 6: Market Position (EMPTY CONTENT — task target) ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide6, "Market Position", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    # Intentionally NO chart or data content — the agent task is to add a pie chart here

    # ---- Slide 7: Growth Strategy ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide7, "Growth Strategy 2025", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    add_bullet_list(slide7, [
        "Expand enterprise sales team by 30% with focus on EMEA and APAC regions",
        "Launch AI-powered analytics suite targeting mid-market segment ($5M pipeline)",
        "Deepen Meridian Corp. partnership with co-developed integration modules",
        "Achieve ISO 27001 certification to unlock government sector opportunities",
        "Increase annual recurring revenue (ARR) target to $185M by end of FY2025",
        "Establish developer ecosystem program with $2M incentive fund for API partners",
    ], Inches(0.8), Inches(1.5), Inches(11), Inches(5), font_size=16, color=DARK_TEXT)

    # ---- Slide 8: Financial Projections ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide8, "Financial Projections", Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                      font_size=32, color=ACCENT_BLUE)
    table_shape2 = slide8.shapes.add_table(5, 3, Inches(1.5), Inches(1.5), Inches(8), Inches(3))
    table2 = table_shape2.table
    headers2 = ["Metric", "FY2024 (Actual)", "FY2025 (Projected)"]
    proj_data = [
        ["Revenue", "$142M", "$185M"],
        ["Gross Margin", "72%", "75%"],
        ["Operating Income", "$28.4M", "$42.5M"],
        ["Customer Count", "2,847", "3,600+"],
    ]
    for c, h in enumerate(headers2):
        cell = table2.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '0072C6'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)
    for r, row_data in enumerate(proj_data, 1):
        for c, val in enumerate(row_data):
            table2.cell(r, c).text = val

    # ---- Slide 9: Contact & Q&A ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    fill9 = slide9.background.fill
    fill9.solid()
    fill9.fore_color.rgb = DARK_BG
    add_title_textbox(slide9, "Thank You", Inches(1), Inches(2), Inches(11), Inches(1.2),
                      font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_body_textbox(slide9, "Questions & Discussion", Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                     font_size=24, color=RGBColor(0xAA, 0xAA, 0xCC))
    add_body_textbox(slide9, "Elena Rodriguez  |  elena.rodriguez@company.com  |  +1 (415) 555-0198",
                     Inches(1), Inches(5), Inches(11), Inches(0.5),
                     font_size=16, color=RGBColor(0x88, 0x88, 0xAA))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
