"""
Reward Script: Create a pie chart on slide 6 with market share data
Task ID: impress_sales_027
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Pie chart exists on slide 6
  Component 2 (0.25): Chart title is 'Market Share 2024'
  Component 3 (0.25): Correct category labels
  Component 4 (0.25): Correct percentage values
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_027'


def persist_app_state():
    """Save any unsaved LibreOffice Impress edits before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        import time
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"PRECONDITION FAIL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # 0-indexed, slide 6

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Pie chart exists on slide 6 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Verify it is specifically a PIE chart (type 5)
            if chart.chart_type == XL_CHART_TYPE.PIE:
                print(f"PASS: Component 1 — Pie chart found on slide 6 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart.chart_type}, expected PIE (5)")
        else:
            print("FAIL: Component 1 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Components 2-4 require a chart to exist
    if chart_shape is None:
        print("FAIL: Components 2-4 skipped — no chart on slide 6")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Market Share 2024' (0.25 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Market Share 2024':
                print(f"PASS: Component 2 — Chart title is 'Market Share 2024' (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Chart title is '{title_text}', expected 'Market Share 2024'")
        else:
            print("FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct category labels (0.25 points)
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        expected_cats = ['Our Product', 'Competitor A', 'Competitor B', 'Others']
        # Normalize: strip whitespace for comparison
        cats_normalized = [c.strip() for c in categories]
        expected_normalized = [c.strip() for c in expected_cats]
        if cats_normalized == expected_normalized:
            print(f"PASS: Component 3 — Categories match: {categories} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Categories are {categories}, expected {expected_cats}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct percentage values (35, 25, 20, 20) (0.25 points)
    try:
        plot = chart.plots[0]
        if len(plot.series) > 0:
            values = list(plot.series[0].values)
            expected_values = [35.0, 25.0, 20.0, 20.0]
            # Allow small floating point tolerance
            if len(values) == len(expected_values):
                all_match = all(
                    abs(v - e) < 0.01
                    for v, e in zip(values, expected_values)
                )
                if all_match:
                    print(f"PASS: Component 4 — Values match: {values} (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 4 — Values are {values}, expected {expected_values}")
            else:
                print(f"FAIL: Component 4 — {len(values)} data points, expected {len(expected_values)}")
        else:
            print("FAIL: Component 4 — No series found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
