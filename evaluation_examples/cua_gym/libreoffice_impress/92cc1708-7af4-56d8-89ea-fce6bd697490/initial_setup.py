"""
Initial Setup: Add a scatter chart on slide 6 of Study_Habits presentation
Task ID: impress_teach_061
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_061'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:  # body placeholder
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(body_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Study Habits & Academic Performance"
    slide1.placeholders[1].text = "A Comprehensive Analysis of Student Learning Patterns\nResearch Division - Spring 2025"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Introduction", [
        "Understanding the relationship between study habits and academic outcomes",
        "has been a central focus of educational research for decades.",
        "This presentation examines data collected from 450 undergraduate",
        "students across three universities during the 2024-2025 academic year.",
        "Our goal is to identify patterns that can inform better study strategies.",
    ])

    # --- Slide 3: Survey Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Survey Methodology", [
        "Online questionnaire distributed to 450 students",
        "Self-reported daily study hours tracked over 12 weeks",
        "Test scores collected from midterm and final exams",
        "Controlled for major, year, and prior GPA",
        "Statistical analysis performed using SPSS and Python",
    ])

    # --- Slide 4: Demographics ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Participant Demographics", [
        "52% female, 46% male, 2% non-binary",
        "Age range: 18-24 (mean: 20.3 years)",
        "Freshmen: 28%, Sophomores: 25%, Juniors: 24%, Seniors: 23%",
        "STEM majors: 40%, Humanities: 30%, Social Sciences: 20%, Other: 10%",
        "Average self-reported GPA prior to study: 3.12",
    ])

    # --- Slide 5: Key Findings ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Key Findings", [
        "Strong positive correlation (r=0.89) between study hours and scores",
        "Diminishing returns observed beyond 9 hours of daily study",
        "Consistent study schedules outperformed irregular cramming by 15%",
        "Group study sessions boosted retention by an average of 12%",
        "Students who took regular breaks scored 8% higher on average",
    ])

    # --- Slide 6: Correlation Analysis (NO CHART - task target) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box
    txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Correlation Analysis"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # --- Slide 7: Recommendations ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Recommendations", [
        "Aim for 5-7 hours of focused study per day",
        "Use spaced repetition techniques for long-term retention",
        "Incorporate active recall through practice tests",
        "Schedule study blocks with 10-minute breaks every 50 minutes",
        "Form study groups of 3-4 students for collaborative learning",
    ])

    # --- Slide 8: Conclusion ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Conclusion", [
        "Study hours strongly predict test performance up to a threshold",
        "Quality of study time matters as much as quantity",
        "Future research should explore subject-specific study strategies",
        "Institutional support programs should promote evidence-based habits",
        "Thank you for your attention - Questions welcome",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
