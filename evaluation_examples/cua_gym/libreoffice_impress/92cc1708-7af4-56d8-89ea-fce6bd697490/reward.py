"""
Reward Script: Add scatter chart on slide 6 with study hours vs test scores
Task ID: impress_teach_061
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Scatter chart exists on slide 6
  Component 2 (0.15): Chart title is 'Study Hours vs Test Score'
  Component 3 (0.20): Axis labels correct (X='Hours Studied', Y='Score (%)')
  Component 4 (0.40): Correct 7 data points matching specified values
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_061'

# Expected data
EXPECTED_X = [2.0, 3.0, 4.0, 5.0, 6.0, 7.0, 8.0]
EXPECTED_Y = [55.0, 62.0, 70.0, 75.0, 82.0, 88.0, 91.0]


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for", domain)
    except Exception as e:
        print("PERSIST_WARN: save hook failed:", e)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # slide 6 (0-indexed)

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Scatter chart exists on slide 6 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            # Check it's a scatter type (XY_SCATTER family)
            ct = chart.chart_type
            # XY_SCATTER = -4169, XY_SCATTER_LINES = 74, XY_SCATTER_LINES_NO_MARKERS = 75,
            # XY_SCATTER_SMOOTH = 72, XY_SCATTER_SMOOTH_NO_MARKERS = 73
            scatter_types = {-4169, 72, 73, 74, 75}
            if int(ct) in scatter_types:
                print(f"PASS: Component 1 — Scatter chart found on slide 6, type={ct} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but not scatter type, got {ct}")
        else:
            print("FAIL: Component 1 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Study Hours vs Test Score' (0.15 points)
    try:
        if chart.chart_title and chart.chart_title.has_text_frame:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'Study Hours vs Test Score':
                print(f"PASS: Component 2 — Chart title matches exactly (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 — Chart title is '{title_text}', expected 'Study Hours vs Test Score'")
        else:
            print("FAIL: Component 2 — Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Axis labels correct (0.20 points)
    # X-axis = 'Hours Studied' (0.10), Y-axis = 'Score (%)' (0.10)
    try:
        axis_score = 0.0

        # For scatter charts: category_axis = X axis, value_axis = Y axis
        try:
            cat_axis = chart.category_axis
            if cat_axis.has_title:
                x_label = cat_axis.axis_title.text_frame.text.strip()
                if x_label == 'Hours Studied':
                    print(f"PASS: Component 3a — X-axis label is 'Hours Studied' (0.10 pts)")
                    axis_score += 0.10
                else:
                    print(f"FAIL: Component 3a — X-axis label is '{x_label}', expected 'Hours Studied'")
            else:
                print("FAIL: Component 3a — X-axis has no title")
        except Exception as e:
            print(f"FAIL: Component 3a — Cannot access X-axis: {e}")

        try:
            val_axis = chart.value_axis
            if val_axis.has_title:
                y_label = val_axis.axis_title.text_frame.text.strip()
                if y_label == 'Score (%)':
                    print(f"PASS: Component 3b — Y-axis label is 'Score (%)' (0.10 pts)")
                    axis_score += 0.10
                else:
                    print(f"FAIL: Component 3b — Y-axis label is '{y_label}', expected 'Score (%)'")
            else:
                print("FAIL: Component 3b — Y-axis has no title")
        except Exception as e:
            print(f"FAIL: Component 3b — Cannot access Y-axis: {e}")

        if axis_score > 0:
            total_score += axis_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct data points (0.40 points)
    # Verify via XML parsing for precise X and Y values
    try:
        ns = {
            'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
        }

        x_vals = []
        y_vals = []

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [f for f in zf.namelist() if 'chart' in f.lower() and f.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.fromstring(f.read())
                    scatter = root.find('.//c:scatterChart', ns)
                    if scatter is not None:
                        for ser in scatter.findall('.//c:ser', ns):
                            xVal = ser.find('c:xVal', ns)
                            if xVal is not None:
                                numCache = xVal.find('.//c:numCache', ns)
                                if numCache is not None:
                                    for pt in numCache.findall('c:pt', ns):
                                        v = pt.find('c:v', ns)
                                        if v is not None:
                                            x_vals.append(float(v.text))

                            yVal = ser.find('c:yVal', ns)
                            if yVal is not None:
                                numCache = yVal.find('.//c:numCache', ns)
                                if numCache is not None:
                                    for pt in numCache.findall('c:pt', ns):
                                        v = pt.find('c:v', ns)
                                        if v is not None:
                                            y_vals.append(float(v.text))

        data_score = 0.0

        # Check number of data points (0.10)
        if len(x_vals) == 7 and len(y_vals) == 7:
            print(f"PASS: Component 4a — Correct number of data points: 7 (0.10 pts)")
            data_score += 0.10
        else:
            print(f"FAIL: Component 4a — Expected 7 data points, found X={len(x_vals)}, Y={len(y_vals)}")

        # Check X values match (0.15)
        if x_vals == EXPECTED_X:
            print(f"PASS: Component 4b — X values match exactly: {x_vals} (0.15 pts)")
            data_score += 0.15
        else:
            print(f"FAIL: Component 4b — X values {x_vals} != expected {EXPECTED_X}")

        # Check Y values match (0.15)
        if y_vals == EXPECTED_Y:
            print(f"PASS: Component 4c — Y values match exactly: {y_vals} (0.15 pts)")
            data_score += 0.15
        else:
            print(f"FAIL: Component 4c — Y values {y_vals} != expected {EXPECTED_Y}")

        if data_score > 0:
            total_score += data_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
