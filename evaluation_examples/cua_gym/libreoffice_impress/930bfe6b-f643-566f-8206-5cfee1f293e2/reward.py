"""
Reward Script: Add content placeholder to master slide for slide 2
Task ID: impress_el_066
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Content placeholder exists on the Title Only layout
  Component 2 (0.3): Placeholder position matches x=2cm, y=5cm
  Component 3 (0.3): Placeholder dimensions match 20cm x 12cm
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_el_066'
FILE_PATH = os.path.join(WORKDIR, f'{TASK_ID}.pptx')

# Expected values in EMU (1 cm = 360000 EMU)
EXPECTED_LEFT = 720000    # 2 cm
EXPECTED_TOP = 1800000    # 5 cm
EXPECTED_WIDTH = 7200000  # 20 cm
EXPECTED_HEIGHT = 4320000 # 12 cm
TOLERANCE = 0.05  # 5% relative tolerance for position/size


def is_approx(actual, expected, tol=TOLERANCE):
    """Check if actual is within tolerance of expected."""
    if expected == 0:
        return actual == 0
    return abs(actual - expected) / abs(expected) <= tol


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that a content placeholder was added to the Title Only slide layout
    (used by slide 2) at position x=2cm, y=5cm with size 20cm x 12cm.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find the layout used by slide 2 (index 1)
    try:
        slide2 = prs.slides[1]
        layout = slide2.slide_layout
        print(f"INFO: Slide 2 uses layout '{layout.name}'")
    except Exception as e:
        print(f"CRITICAL: Cannot access slide 2 or its layout: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find a content/object placeholder on the layout (not Title, Date, Footer, SlideNum)
    # In the initial state, the Title Only layout has only Title + date/footer/slidenum
    # The task adds a content placeholder (OBJECT type, typically idx=1)
    # Filter by TYPE not index, since placeholder indices vary across VM environments
    content_ph = None
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
        # Content placeholder types: OBJECT (7) and BODY (2)
        CONTENT_TYPES = {7, 2}  # PP_PLACEHOLDER_TYPE.OBJECT, PP_PLACEHOLDER_TYPE.BODY
        for ph in layout.placeholders:
            ph_type_val = int(ph.placeholder_format.type)
            # Only accept OBJECT or BODY type placeholders as content placeholders
            if ph_type_val in CONTENT_TYPES:
                content_ph = ph
                print(f"INFO: Found content placeholder: idx={ph.placeholder_format.idx}, "
                      f"type={ph.placeholder_format.type}, name='{ph.name}'")
                print(f"INFO: Position: left={ph.left}, top={ph.top}, "
                      f"width={ph.width}, height={ph.height}")
                break
    except Exception as e:
        print(f"ERROR: Error scanning layout placeholders: {e}")

    # Component 1: Content placeholder exists on the layout (0.4 points)
    try:
        if content_ph is not None:
            print(f"PASS: Component 1 -- Content placeholder found on layout (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 -- No content placeholder found on the slide 2 layout. "
                  f"Only title/date/footer/slidenum placeholders exist.")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Position is correct x=2cm (720000 EMU), y=5cm (1800000 EMU) (0.3 points)
    try:
        if content_ph is not None:
            left_ok = is_approx(content_ph.left, EXPECTED_LEFT)
            top_ok = is_approx(content_ph.top, EXPECTED_TOP)
            if left_ok and top_ok:
                print(f"PASS: Component 2 -- Position correct: left={content_ph.left} "
                      f"(expected {EXPECTED_LEFT}), top={content_ph.top} "
                      f"(expected {EXPECTED_TOP}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 -- Position mismatch: "
                      f"left={content_ph.left} (expected {EXPECTED_LEFT}, ok={left_ok}), "
                      f"top={content_ph.top} (expected {EXPECTED_TOP}, ok={top_ok})")
        else:
            print(f"FAIL: Component 2 -- No content placeholder to check position")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Dimensions correct width=20cm (7200000 EMU), height=12cm (4320000 EMU) (0.3 points)
    try:
        if content_ph is not None:
            width_ok = is_approx(content_ph.width, EXPECTED_WIDTH)
            height_ok = is_approx(content_ph.height, EXPECTED_HEIGHT)
            if width_ok and height_ok:
                print(f"PASS: Component 3 -- Dimensions correct: width={content_ph.width} "
                      f"(expected {EXPECTED_WIDTH}), height={content_ph.height} "
                      f"(expected {EXPECTED_HEIGHT}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 -- Dimensions mismatch: "
                      f"width={content_ph.width} (expected {EXPECTED_WIDTH}, ok={width_ok}), "
                      f"height={content_ph.height} (expected {EXPECTED_HEIGHT}, ok={height_ok})")
        else:
            print(f"FAIL: Component 3 -- No content placeholder to check dimensions")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
