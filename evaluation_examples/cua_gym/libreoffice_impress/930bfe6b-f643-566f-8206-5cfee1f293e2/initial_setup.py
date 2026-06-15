"""
Initial Setup: Create a presentation where slide 2 uses a layout with only a title placeholder.
Task ID: impress_el_066
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_066'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide (layout 0) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Marketing Strategy"
    slide1.placeholders[1].text = "Prepared by Sarah Chen, VP of Marketing\nAcme Corporation"

    # --- Slide 2: Title Only layout (layout 5 = Blank or 6 = Title Only) ---
    # Use Title Only layout (index 5 in default template is often Title Only)
    # We need a layout that has ONLY a title placeholder
    # In default pptx template: layout index 5 = Title Only
    title_only_layout = prs.slide_layouts[5]  # Title Only
    slide2 = prs.slides.add_slide(title_only_layout)
    slide2.shapes.title.text = "Regional Sales Performance"

    # Add a text box with some context (not a placeholder - just a free textbox)
    from pptx.util import Cm as CmUtil
    txBox = slide2.shapes.add_textbox(Cm(2), Cm(5), Cm(8), Cm(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Note: Content area to be configured via master slide editor"
    if p.runs:
        p.runs[0].font.size = Pt(11)
        p.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
        p.runs[0].font.italic = True

    # --- Slide 3: Content slide (layout 1 = Title and Content) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Metrics Overview"
    body = slide3.placeholders[1]
    tf3 = body.text_frame
    tf3.text = "Total Revenue: $4.2M (+12% YoY)"
    p2 = tf3.add_paragraph()
    p2.text = "New Customers: 1,847 (+23% QoQ)"
    p2.level = 0
    p3 = tf3.add_paragraph()
    p3.text = "Customer Retention Rate: 94.3%"
    p3.level = 0
    p4 = tf3.add_paragraph()
    p4.text = "Average Deal Size: $28,500"
    p4.level = 0
    p5 = tf3.add_paragraph()
    p5.text = "Pipeline Value: $12.8M"
    p5.level = 0

    # --- Slide 4: Another Title Only slide ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Product Roadmap Timeline"

    # Add a simple table
    table_shape = slide4.shapes.add_table(4, 3, Cm(2), Cm(5), Cm(22), Cm(8))
    table = table_shape.table
    # Headers
    headers = ["Milestone", "Target Date", "Status"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        for run in table.cell(0, i).text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    # Data
    data = [
        ["Beta Release v3.0", "June 15, 2025", "On Track"],
        ["Enterprise Integration", "August 1, 2025", "Planning"],
        ["Mobile App Launch", "October 20, 2025", "Design Phase"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 5: Closing slide ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Thank You"
    slide5.placeholders[1].text = "Questions? Contact: s.chen@acme-corp.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
