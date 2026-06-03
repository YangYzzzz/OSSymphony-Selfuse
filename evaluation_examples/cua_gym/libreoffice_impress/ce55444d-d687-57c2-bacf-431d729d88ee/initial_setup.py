"""
Initial Setup: 5-slide team onboarding presentation
Task ID: osworld_impress_multi_op_combined_012
Domain: libreoffice_impress
State: Slide 2 title NOT underlined; slide 4 has WHITE background
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen 10x7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Welcome to Northbridge Technologies"
    slide1.placeholders[1].text = "New Employee Onboarding Program\nHR & People Operations"
    # Set white background for slide 1
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 2: Team Introduction (title NOT underlined) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Meet Your Team"
    # Ensure title text is NOT underlined
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.underline = False
            run.font.bold = True
            run.font.size = Pt(36)
    # Content placeholder
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Engineering Team"
    paras = [
        "Alice Morgan – Lead Software Engineer",
        "David Park – Backend Developer",
        "Priya Nair – Frontend Developer",
        "James Liu – DevOps Engineer",
    ]
    for item in paras:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1
    # White background for slide 2
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 3: Onboarding Process ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Your Onboarding Journey"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Week 1–2: Orientation & Setup"
    steps = [
        "Complete IT setup and security training",
        "Meet with your manager for goals alignment",
        "Shadow team members across departments",
        "Week 3–4: Active Contribution",
        "Join sprint planning and standups",
        "Complete your first feature or deliverable",
    ]
    for step in steps:
        p = tf3.add_paragraph()
        p.text = step
        p.level = 1
    # White background for slide 3
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 4: Company Benefits (WHITE background — task target) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Benefits & Perks"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Health & Wellness"
    benefits = [
        "Comprehensive medical, dental, and vision coverage",
        "Annual wellness stipend of $1,200",
        "Mental health support and EAP services",
        "Professional Development",
        "Learning & development budget of $2,500/year",
        "Conference attendance and certification reimbursement",
    ]
    for b in benefits:
        p = tf4.add_paragraph()
        p.text = b
        p.level = 1
    # WHITE background for slide 4 (task: change to pale green)
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 5: Next Steps & Contacts ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Key Contacts"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Important Contacts"
    contacts = [
        "HR Business Partner: Sandra Lee (sandra.lee@northbridge.com)",
        "IT Helpdesk: it-support@northbridge.com | ext. 5100",
        "Office Manager: Tom Reyes (tom.reyes@northbridge.com)",
        "What to Do This Week",
        "Review the employee handbook in the HR portal",
        "Set up your profile on the internal directory",
        "Schedule your 30-day check-in with your manager",
    ]
    for c in contacts:
        p = tf5.add_paragraph()
        p.text = c
        p.level = 1
    # White background for slide 5
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
