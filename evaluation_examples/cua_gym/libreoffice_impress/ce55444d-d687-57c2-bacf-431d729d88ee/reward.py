"""
Reward Script: Underline slide 2 title and change slide 4 background to pale green
Task ID: osworld_impress_multi_op_combined_012
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 title text is underlined          — 0.5 points
  Component 2: Slide 4 background is pale green (90EE90) — 0.5 points
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_012'

# Pale green: CSS named color "palegreen" = RGB(144, 238, 144) = hex 90EE90
# Accept a range of pale green shades (tolerance approach)
PALE_GREEN_HEX = '90EE90'


def is_pale_green(rgb_str):
    """
    Returns True if the given RGB hex string is a pale green color.
    Accepts the canonical pale green (90EE90) and similar light greens.
    Tolerance: each channel within ±30 of the reference value.
    Reference: R=144, G=238, B=144
    """
    try:
        r = int(rgb_str[0:2], 16)
        g = int(rgb_str[2:4], 16)
        b = int(rgb_str[4:6], 16)
        # Must be a greenish color: G dominant, R and B close to each other and lower than G
        # Accept broad "pale green" range
        ref_r, ref_g, ref_b = 144, 238, 144
        tolerance = 35
        return (abs(r - ref_r) <= tolerance and
                abs(g - ref_g) <= tolerance and
                abs(b - ref_b) <= tolerance)
    except Exception:
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Component 1 (0.5): Slide 2 title run(s) have underline=True
    Component 2 (0.5): Slide 4 background fill is solid and pale green
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 title is underlined (0.5 points)
    # The task says to underline the title on slide 2.
    # Initial: underline=False; Golden: underline=True
    try:
        slide2 = prs.slides[1]  # 0-indexed → slide 2
        title_shape = None
        for shape in slide2.shapes:
            if shape.has_text_frame:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    # TITLE placeholder type = 1 (CENTER_TITLE=3 also considered)
                    if ph_type in (1, 3):
                        title_shape = shape
                        break
        if title_shape is None:
            print("FAIL: Component 1 — No title placeholder found on slide 2")
        else:
            title_text = title_shape.text_frame.text.strip()
            # Check that at least one non-empty run has underline=True
            underlined_runs = 0
            total_runs = 0
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or '').strip():
                        total_runs += 1
                        if run.font.underline is True:
                            underlined_runs += 1
            if total_runs > 0 and underlined_runs == total_runs:
                print(f"PASS: Component 1 — Slide 2 title '{title_text}' is fully underlined ({underlined_runs}/{total_runs} runs) (0.5 pts)")
                total_score += 0.5
            elif underlined_runs > 0:
                print(f"FAIL: Component 1 — Slide 2 title '{title_text}' partially underlined ({underlined_runs}/{total_runs} runs); need all runs underlined")
            else:
                print(f"FAIL: Component 1 — Slide 2 title '{title_text}' is NOT underlined (underline={[r.font.underline for r in para.runs for para in title_shape.text_frame.paragraphs if (r.text or '').strip()]})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 4 background is pale green (0.5 points)
    # The task says to change slide 4 background to pale green.
    # Initial: FFFFFF (white); Golden: 90EE90 (pale green)
    try:
        slide4 = prs.slides[3]  # 0-indexed → slide 4
        fill = slide4.background.fill
        if fill.type is None:
            print("FAIL: Component 2 — Slide 4 background has no fill (not solid)")
        elif fill.type == 1:  # SOLID fill
            try:
                rgb = fill.fore_color.rgb
                rgb_str = str(rgb)  # e.g. "90EE90"
                if is_pale_green(rgb_str):
                    print(f"PASS: Component 2 — Slide 4 background is pale green (#{rgb_str}) (0.5 pts)")
                    total_score += 0.5
                else:
                    print(f"FAIL: Component 2 — Slide 4 background color is #{rgb_str}, expected pale green (~#90EE90)")
            except Exception as ce:
                print(f"FAIL: Component 2 — Could not read slide 4 background RGB: {ce}")
        else:
            print(f"FAIL: Component 2 — Slide 4 background fill type is {fill.type}, expected SOLID (1) with pale green color")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
