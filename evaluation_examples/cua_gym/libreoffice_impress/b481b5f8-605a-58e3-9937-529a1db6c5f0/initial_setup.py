"""
Initial Setup: Cybersecurity Awareness Training presentation (15 slides)
Task ID: impress_wf_026
Domain: libreoffice_impress

Creates a 15-slide cybersecurity training module with:
- Title slide, 7 lesson slides, 5 quiz slides (no animations), scoring rubric, certificate
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_points(text_frame, bullets, font_size=14, color=None):
    """Add bullet point paragraphs to an existing text frame."""
    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            if color:
                run.font.color.rgb = color


def create_lesson_slide(prs, title_text, bullets):
    """Create a lesson content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Dark blue header bar background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    bg_shape.line.fill.background()

    # Title text
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.9),
                title_text, font_size=28, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)

    # Bullet content area
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    add_bullet_points(tf, bullets, font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    return slide


def create_quiz_slide(prs, question, options, option_colors):
    """Create a quiz slide with question and 4 colored rectangle answer options."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Question header background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    bg_shape.line.fill.background()

    # Question text
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.9),
                question, font_size=22, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)

    # Four answer rectangles (A-D) arranged in 2x2 grid
    labels = ['A', 'B', 'C', 'D']
    positions = [
        (Inches(0.5), Inches(1.8)),   # A - top left
        (Inches(5.2), Inches(1.8)),   # B - top right
        (Inches(0.5), Inches(4.0)),   # C - bottom left
        (Inches(5.2), Inches(4.0)),   # D - bottom right
    ]
    rect_width = Inches(4.3)
    rect_height = Inches(1.6)

    for i, (label, option_text) in enumerate(zip(labels, options)):
        left, top = positions[i]
        color = option_colors[i]

        # Colored rectangle
        rect = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, rect_width, rect_height
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        rect.line.width = Pt(1.5)

        # Add text inside rectangle
        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{label}. {option_text}"
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Dark background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # Shield icon shape
    lock_icon = slide1.shapes.add_shape(
        MSO_SHAPE.PENTAGON, Inches(4.0), Inches(0.8), Inches(2), Inches(2)
    )
    lock_icon.fill.solid()
    lock_icon.fill.fore_color.rgb = RGBColor(0x00, 0x7B, 0xFF)
    lock_icon.line.fill.background()

    add_textbox(slide1, Inches(1), Inches(3.2), Inches(8), Inches(1.2),
                "Cybersecurity Awareness Training", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, Inches(1.5), Inches(4.6), Inches(7), Inches(0.8),
                "Protecting Our Digital Workplace — 2025 Edition", font_size=18,
                bold=False, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

    add_textbox(slide1, Inches(2), Inches(5.8), Inches(6), Inches(0.6),
                "Information Security Department | Meridian Technologies", font_size=14,
                bold=False, color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

    # ---- Slides 2-8: Lesson Content ----
    lessons = [
        ("Password Security", [
            "Use passwords with at least 12 characters combining letters, numbers, and symbols",
            "Never reuse passwords across different accounts or services",
            "Enable multi-factor authentication (MFA) on all critical accounts",
            "Use a reputable password manager like LastPass or 1Password",
            "Change passwords immediately if a breach is suspected",
            "Avoid using personal information such as birthdays or pet names",
        ]),
        ("Phishing Awareness", [
            "Verify sender email addresses carefully — look for subtle misspellings",
            "Hover over links before clicking to preview the actual destination URL",
            "Be suspicious of urgent requests for personal information or payments",
            "Report suspected phishing emails to security@meridiantech.com immediately",
            "Never download attachments from unknown or unexpected senders",
            "Watch for grammatical errors and unusual formatting in emails",
        ]),
        ("Malware Prevention", [
            "Keep operating systems and software updated with the latest patches",
            "Only download software from official sources and verified repositories",
            "Run scheduled antivirus scans at least weekly on all work devices",
            "Disable auto-run for USB drives and external media",
            "Back up critical files to approved cloud storage regularly",
            "Report any unusual system behavior to IT support within 1 hour",
        ]),
        ("Social Engineering", [
            "Verify identity before sharing sensitive information — even with colleagues",
            "Be cautious of unsolicited phone calls claiming to be from IT support",
            "Never share access badges or allow tailgating through secure doors",
            "Question unusual requests even if they appear to come from leadership",
            "Document and report any social engineering attempts to the security team",
            "Remember: legitimate IT staff will never ask for your password",
        ]),
        ("Data Protection", [
            "Classify all documents according to company data sensitivity levels",
            "Encrypt sensitive files before transmitting via email or file sharing",
            "Lock your workstation (Win+L or Ctrl+Cmd+Q) when stepping away",
            "Store confidential documents only in approved secure locations",
            "Shred physical documents containing sensitive information",
            "Follow the data retention policy — delete files past their retention date",
        ]),
        ("Secure Browsing", [
            "Look for HTTPS and the padlock icon before entering credentials",
            "Avoid using public Wi-Fi for accessing company resources",
            "Use the corporate VPN when working remotely or traveling",
            "Clear browser cache and cookies after using shared computers",
            "Do not install unauthorized browser extensions or plugins",
            "Bookmark frequently used internal sites to avoid URL typos",
        ]),
        ("Incident Reporting", [
            "Report security incidents within 30 minutes of discovery",
            "Use the Security Incident Report Form on the company intranet",
            "Preserve evidence — do not delete suspicious emails or files",
            "Contact the Security Operations Center (SOC) at ext. 4500 for urgent issues",
            "Document the timeline: when you noticed it, what you observed, actions taken",
            "Follow up with the assigned incident handler within 24 hours",
        ]),
    ]

    for title, bullets in lessons:
        create_lesson_slide(prs, title, bullets)

    # ---- Slides 9-13: Quiz Slides (NO animations in initial) ----
    quiz_data = [
        {
            "question": "Q1: What is the minimum recommended password length?",
            "options": ["6 characters", "8 characters", "12 characters", "16 characters"],
        },
        {
            "question": "Q2: What should you do when you receive a suspicious email?",
            "options": [
                "Open attachments to investigate",
                "Forward it to all colleagues",
                "Report it to the security team",
                "Delete it without telling anyone",
            ],
        },
        {
            "question": "Q3: Which is the safest way to download software?",
            "options": [
                "From email attachments",
                "From official app stores",
                "From social media links",
                "From pop-up advertisements",
            ],
        },
        {
            "question": "Q4: What should you do before sharing sensitive information?",
            "options": [
                "Post it on the company chat",
                "Verify the requester's identity",
                "Send it via personal email",
                "Share it on social media",
            ],
        },
        {
            "question": "Q5: How quickly should security incidents be reported?",
            "options": [
                "Within one week",
                "Within 24 hours",
                "Within 30 minutes",
                "Only if it causes damage",
            ],
        },
    ]

    quiz_colors = [
        RGBColor(0xE7, 0x4C, 0x3C),  # Red
        RGBColor(0x34, 0x98, 0xDB),  # Blue
        RGBColor(0x2E, 0xCC, 0x71),  # Green
        RGBColor(0xF3, 0x9C, 0x12),  # Orange
    ]

    for q in quiz_data:
        create_quiz_slide(prs, q["question"], q["options"], quiz_colors)

    # ---- Slide 14: Scoring Rubric ----
    slide14 = prs.slides.add_slide(prs.slide_layouts[5])

    # Header
    bg_shape = slide14.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    bg_shape.line.fill.background()

    add_textbox(slide14, Inches(0.5), Inches(0.15), Inches(9), Inches(0.9),
                "Quiz Scoring Rubric", font_size=28, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)

    # Scoring table: 4 rows (header + 4 score ranges)
    table_shape = slide14.shapes.add_table(
        5, 3, Inches(1.5), Inches(1.8), Inches(7), Inches(3.5)
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)

    # Header row
    headers = ["Score Range", "Rating", "Recommended Action"]
    header_color = RGBColor(0x1B, 0x3A, 0x5C)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # Data rows
    rubric_data = [
        ["0-1 correct", "Needs Improvement", "Mandatory retraining within 1 week"],
        ["2-3 correct", "Fair", "Review weak areas and retake quiz"],
        ["4 correct", "Good", "Annual refresher recommended"],
        ["5 correct", "Excellent", "No further action required"],
    ]
    row_colors = [
        RGBColor(0xFD, 0xED, 0xEC),
        RGBColor(0xFE, 0xF9, 0xE7),
        RGBColor(0xE8, 0xF8, 0xF5),
        RGBColor(0xE8, 0xF6, 0xEF),
    ]
    for row_idx, (row_data, row_bg) in enumerate(zip(rubric_data, row_colors), 1):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
                run.font.name = "Arial"
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg

    # ---- Slide 15: Certificate Template ----
    slide15 = prs.slides.add_slide(prs.slide_layouts[5])

    # Background
    fill = slide15.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFD, 0xF5, 0xE6)

    # Decorative border - outer rectangle
    border_outer = slide15.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(0.3),
        Inches(9.4), Inches(6.9)
    )
    border_outer.fill.background()
    border_outer.line.color.rgb = RGBColor(0xC8, 0x96, 0x3E)
    border_outer.line.width = Pt(4)

    # Inner border
    border_inner = slide15.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6),
        Inches(8.8), Inches(6.3)
    )
    border_inner.fill.background()
    border_inner.line.color.rgb = RGBColor(0xC8, 0x96, 0x3E)
    border_inner.line.width = Pt(2)

    add_textbox(slide15, Inches(1), Inches(1.0), Inches(8), Inches(0.8),
                "CERTIFICATE OF COMPLETION", font_size=30, bold=True,
                color=RGBColor(0xC8, 0x96, 0x3E), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(1), Inches(2.0), Inches(8), Inches(0.6),
                "This certifies that", font_size=16,
                bold=False, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(1.5), Inches(2.8), Inches(7), Inches(0.8),
                "[Employee Name]", font_size=28, bold=True,
                color=RGBColor(0x1B, 0x3A, 0x5C), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(1), Inches(3.8), Inches(8), Inches(0.8),
                "has successfully completed the Cybersecurity Awareness Training\nand demonstrated proficiency in organizational security practices.",
                font_size=14, bold=False,
                color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(1), Inches(5.2), Inches(3.5), Inches(0.5),
                "Date: _______________", font_size=12,
                bold=False, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(5.5), Inches(5.2), Inches(3.5), Inches(0.5),
                "Signature: _______________", font_size=12,
                bold=False, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

    add_textbox(slide15, Inches(2), Inches(6.0), Inches(6), Inches(0.4),
                "Meridian Technologies — Information Security Department", font_size=10,
                bold=False, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
