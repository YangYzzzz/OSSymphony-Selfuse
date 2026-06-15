"""
Reward Script: Master slide with three-zone layout
Task ID: impress_gf2_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Top branding bar: full-width rectangle, 1.5cm tall, color #2563EB
  Component 2 (0.15) - Logo placeholder inside top bar
  Component 3 (0.30) - Bottom footer strip: full-width rectangle, 1.0cm tall, color #1E3A5F
  Component 4 (0.25) - Footer placeholders (slide number, date, company name) in bottom strip
"""

import os
from pptx import Presentation
from pptx.util import Emu, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_010'


def approx_equal(val1, val2, tolerance=0.02):
    """Check if two values are approximately equal within relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 50000  # ~0.14cm tolerance for zero comparisons
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def get_shape_fill_color(shape):
    """Get the solid fill color of a shape as a hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get slide master (first one)
    if len(prs.slide_masters) == 0:
        print("CRITICAL: No slide masters found")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]
    slide_width = prs.slide_width  # expected ~12192119 EMU = 33.87cm

    # Collect auto shapes and placeholders on master
    auto_shapes = []
    placeholders = []
    for shape in master.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            placeholders.append(shape)

    # --- Component 1: Top branding bar (0.30 points) ---
    # Full-width rectangle at top, 1.5cm tall, color #2563EB
    try:
        top_bar = None
        for shape in auto_shapes:
            color = get_shape_fill_color(shape)
            # Look for a shape at the top of the slide with blue fill
            if (color and color.upper() == '2563EB'
                    and shape.top < Cm(1)
                    and approx_equal(shape.width, slide_width)):
                top_bar = shape
                break

        if top_bar is not None:
            height_cm = top_bar.height / 360000
            if approx_equal(top_bar.height, Cm(1.5)):
                print(f"PASS: Component 1 — Top branding bar found: full-width, {height_cm:.2f}cm tall, color #2563EB (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Top bar found but height={height_cm:.2f}cm, expected 1.5cm")
        else:
            print("FAIL: Component 1 — No top branding bar found (expected full-width, #2563EB, at top)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # --- Component 2: Logo placeholder inside top bar (0.15 points) ---
    # A shape within the top bar area with logo-related content
    try:
        logo_found = False
        if top_bar is not None:
            top_bar_bottom = top_bar.top + top_bar.height
            for shape in auto_shapes:
                # Logo placeholder should be within the top bar region
                if (shape != top_bar
                        and shape.top >= top_bar.top
                        and (shape.top + shape.height) <= top_bar_bottom + Cm(0.2)):
                    # Check if it has logo-related text or is a placeholder box
                    shape_text = ''
                    if hasattr(shape, 'text'):
                        shape_text = shape.text.strip().lower()
                    if 'logo' in shape_text or shape_text == '':
                        logo_found = True
                        break

            # Also check for picture shapes or placeholders in the top bar
            if not logo_found:
                for shape in master.shapes:
                    if (shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER)
                            and shape.top >= top_bar.top
                            and (shape.top + shape.height) <= top_bar_bottom + Cm(0.2)
                            and shape not in placeholders[:5]):  # skip original placeholders
                        logo_found = True
                        break

        if logo_found:
            print(f"PASS: Component 2 — Logo placeholder found inside top bar (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 2 — No logo placeholder found inside top branding bar")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # --- Component 3: Bottom footer strip (0.30 points) ---
    # Full-width rectangle at bottom, 1.0cm tall, color #1E3A5F
    try:
        bottom_bar = None
        slide_height = prs.slide_height
        for shape in auto_shapes:
            color = get_shape_fill_color(shape)
            # Look for a shape near the bottom with dark blue fill
            if (color and color.upper() == '1E3A5F'
                    and approx_equal(shape.width, slide_width)
                    and (shape.top + shape.height) >= slide_height - Cm(0.5)):
                bottom_bar = shape
                break

        if bottom_bar is not None:
            height_cm = bottom_bar.height / 360000
            if approx_equal(bottom_bar.height, Cm(1.0)):
                print(f"PASS: Component 3 — Bottom footer strip found: full-width, {height_cm:.2f}cm tall, color #1E3A5F (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 3 — Bottom bar found but height={height_cm:.2f}cm, expected 1.0cm")
        else:
            print("FAIL: Component 3 — No bottom footer strip found (expected full-width, #1E3A5F, at bottom)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # --- Component 4: Footer placeholders in bottom strip (0.25 points) ---
    # Slide number, date, and company name placeholders within the bottom strip
    try:
        if bottom_bar is not None:
            bar_top = bottom_bar.top
            bar_bottom = bottom_bar.top + bottom_bar.height

            footer_items_found = 0
            # Check all shapes (placeholders and auto shapes) in the bottom strip region
            for shape in master.shapes:
                if shape == bottom_bar:
                    continue
                # Shape must overlap with the bottom strip
                shape_bottom = shape.top + shape.height
                if shape.top >= bar_top - Cm(0.2) and shape.top < bar_bottom + Cm(0.2):
                    shape_text = ''
                    if hasattr(shape, 'text'):
                        shape_text = shape.text.strip().lower()
                    # Look for date, slide number, or company name
                    if any(kw in shape.name.lower() for kw in ['date', 'number', 'slide num', 'company', 'footer']):
                        footer_items_found += 1
                    elif any(kw in shape_text for kw in ['date', '#', 'meridian', 'company']):
                        footer_items_found += 1

            # We expect at least 3 footer items (date, company name, slide number)
            if footer_items_found >= 3:
                print(f"PASS: Component 4 — {footer_items_found} footer placeholders found in bottom strip (0.25 pts)")
                total_score += 0.25
            elif footer_items_found >= 2:
                partial = 0.15
                print(f"PARTIAL: Component 4 — {footer_items_found}/3 footer placeholders found ({partial} pts)")
                total_score += partial
            elif footer_items_found >= 1:
                partial = 0.08
                print(f"PARTIAL: Component 4 — {footer_items_found}/3 footer placeholders found ({partial} pts)")
                total_score += partial
            else:
                print("FAIL: Component 4 — No footer placeholders found in bottom strip")
        else:
            print("FAIL: Component 4 — Cannot check footers without bottom bar")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
