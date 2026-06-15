"""
Initial Setup: Corporate presentation with plain master slide
Task ID: impress_gf2_010
Domain: libreoffice_impress

Creates a 15-slide corporate presentation with a plain white master slide
(title + content placeholders only, no decorative elements).
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Standard widescreen 33.867 cm x 19.05 cm (13.333 x 7.5 inches)
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Meridian Technologies"
    slide1.placeholders[1].text = "Annual Strategic Review 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Company Overview & Mission"
    items2 = [
        "Financial Performance Summary",
        "Product Portfolio Update",
        "Market Analysis & Competitive Landscape",
        "Technology Roadmap",
        "Human Resources & Culture",
        "Customer Success Stories",
        "Risk Assessment",
        "Strategic Priorities for 2026",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Company Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Company Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Founded in 2011 by Dr. Elena Vasquez and James Park"
    for txt in [
        "Headquarters: San Francisco, CA with offices in London, Tokyo, and Sydney",
        "Employees: 2,847 across 12 countries",
        "Core Focus: Enterprise cloud infrastructure and AI-driven analytics",
        "Annual Revenue: $428M (FY2024)",
        "Fortune 500 clients: 87 active accounts",
    ]:
        p = body3.add_paragraph()
        p.text = txt
        p.level = 0

    # --- Slide 4: Financial Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Highlights FY2024"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Revenue: $428M (+18% YoY)"
    for txt in [
        "Gross Margin: 72.3% (up from 69.1%)",
        "Operating Income: $94.2M",
        "Free Cash Flow: $67.8M",
        "R&D Investment: $112M (26% of revenue)",
        "Customer Retention Rate: 96.4%",
    ]:
        p = body4.add_paragraph()
        p.text = txt

    # --- Slide 5: Revenue Breakdown ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Revenue by Segment"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Cloud Platform: $198M (46%)"
    for txt in [
        "Analytics Suite: $132M (31%)",
        "Professional Services: $58M (14%)",
        "Support & Maintenance: $40M (9%)",
    ]:
        p = body5.add_paragraph()
        p.text = txt

    # --- Slide 6: Product Portfolio ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Product Portfolio"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "MeridianCloud: Enterprise-grade IaaS platform"
    for txt in [
        "InsightEngine: AI-powered business analytics",
        "DataVault: Secure data lake management",
        "FlowConnect: API integration middleware",
        "GuardShield: Zero-trust security framework",
    ]:
        p = body6.add_paragraph()
        p.text = txt

    # --- Slide 7: Market Analysis ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Market Analysis"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Total Addressable Market: $89B by 2026"
    for txt in [
        "Cloud infrastructure growing at 22% CAGR",
        "AI analytics segment accelerating to 35% CAGR",
        "Key competitors: Cloudium, DataSphere, NexaTech",
        "Meridian market share: 4.8% (up from 3.2% in 2023)",
        "Strongest growth in APAC region (+34% YoY)",
    ]:
        p = body7.add_paragraph()
        p.text = txt

    # --- Slide 8: Technology Roadmap ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Technology Roadmap 2025-2026"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Q1 2025: Launch InsightEngine v3.0 with generative AI"
    for txt in [
        "Q2 2025: MeridianCloud multi-region failover",
        "Q3 2025: DataVault real-time streaming ingestion",
        "Q4 2025: FlowConnect GraphQL support",
        "Q1 2026: GuardShield quantum-resistant encryption",
        "Q2 2026: Unified platform dashboard release",
    ]:
        p = body8.add_paragraph()
        p.text = txt

    # --- Slide 9: HR & Culture ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "People & Culture"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Employee satisfaction score: 4.3/5.0 (Gallup Q12)"
    for txt in [
        "Voluntary turnover: 8.7% (industry avg: 13.2%)",
        "New hires in 2024: 412 across engineering and sales",
        "Diversity initiatives: 42% women in leadership",
        "Learning & development budget: $3.2M",
        "Remote work policy: hybrid (3 days in-office)",
    ]:
        p = body9.add_paragraph()
        p.text = txt

    # --- Slide 10: Customer Success ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Customer Success Stories"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "GlobalBank Corp: 40% reduction in data processing time"
    for txt in [
        "HealthFirst Systems: HIPAA-compliant analytics in 6 weeks",
        "RetailMax: $12M annual savings from cloud migration",
        "AutoDrive Inc: Real-time telemetry for 50K vehicles",
        "EduBright: Personalized learning analytics for 2M students",
    ]:
        p = body10.add_paragraph()
        p.text = txt

    # --- Slide 11: Risk Assessment ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Risk Assessment"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "Regulatory: GDPR/CCPA compliance costs increasing"
    for txt in [
        "Competitive: Large cloud providers expanding analytics offerings",
        "Talent: AI/ML engineer shortage driving salary inflation",
        "Economic: Enterprise spending caution in uncertain macro environment",
        "Technical: Legacy system integration complexity",
    ]:
        p = body11.add_paragraph()
        p.text = txt

    # --- Slide 12: Strategic Priorities ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Strategic Priorities 2026"
    body12 = slide12.placeholders[1].text_frame
    body12.text = "Expand APAC presence with Tokyo data center"
    for txt in [
        "Launch vertical-specific solutions for healthcare and finance",
        "Achieve SOC 2 Type II and FedRAMP certifications",
        "Grow partner ecosystem to 200+ certified integrators",
        "Increase ARR to $550M through upsell and new logos",
    ]:
        p = body12.add_paragraph()
        p.text = txt

    # --- Slide 13: Investment Thesis ---
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Investment Thesis"
    body13 = slide13.placeholders[1].text_frame
    body13.text = "Strong product-market fit in enterprise cloud + AI"
    for txt in [
        "High gross margins with improving operating leverage",
        "Sticky customer base with 96%+ retention",
        "Experienced leadership team with proven track record",
        "Clear path to $1B ARR within 3 years",
    ]:
        p = body13.add_paragraph()
        p.text = txt

    # --- Slide 14: Partnerships ---
    slide14 = prs.slides.add_slide(prs.slide_layouts[1])
    slide14.shapes.title.text = "Key Partnerships"
    body14 = slide14.placeholders[1].text_frame
    body14.text = "AWS: Preferred technology partner for hybrid deployments"
    for txt in [
        "Accenture: Joint go-to-market for financial services",
        "Snowflake: Native integration for data warehousing",
        "Okta: Single sign-on and identity management",
        "Deloitte: Advisory services for digital transformation clients",
    ]:
        p = body14.add_paragraph()
        p.text = txt

    # --- Slide 15: Thank You / Q&A ---
    slide15 = prs.slides.add_slide(prs.slide_layouts[0])
    slide15.shapes.title.text = "Thank You"
    slide15.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
