"""
Reward Script: Add dark blue header bar to master slide
Task ID: impress_fix_082
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Master slide contains a non-placeholder AUTO_SHAPE rectangle
  Component 2 (0.3): Rectangle positioned at (0,0) with dimensions ~10in x 0.75in
  Component 3 (0.4): Rectangle has solid fill color #003366
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_082'

# Expected values in EMU
EXPECTED_LEFT = 0
EXPECTED_TOP = 0
EXPECTED_WIDTH = Inches(10)   # 9144000 EMU
EXPECTED_HEIGHT = Inches(0.75) # 685800 EMU
EXPECTED_COLOR = '003366'
POSITION_TOLERANCE = 0.02  # 2% relative tolerance


def is_approx(actual, expected, tolerance=POSITION_TOLERANCE):
    """Check if actual is approximately equal to expected."""
    if expected == 0:
        # For zero expected, use absolute tolerance based on slide width
        return abs(actual) <= Inches(0.05)  # within 0.05 inches of 0
    return abs(actual - expected) / max(abs(actual), abs(expected)) <= tolerance


def find_header_bar_shape(master):
    """
    Find a non-placeholder AUTO_SHAPE on the master slide that could be the header bar.
    Returns the shape or None.
    """
    for shape in master.shapes:
        # Skip placeholders -- those are pre-existing
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        # Look for AUTO_SHAPE (rectangles are auto shapes)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file has at least one slide master
    if len(prs.slide_masters) == 0:
        print("FAIL: No slide masters found in presentation")
        print("REWARD: 0.0")
        return 0.0

    master = prs.slide_masters[0]

    # Component 1: Master slide has a non-placeholder AUTO_SHAPE (0.3 points)
    # This shape does NOT exist in the initial file (only 5 placeholders).
    try:
        header_shape = find_header_bar_shape(master)
        if header_shape is not None:
            print(f"PASS: Component 1 -- Found non-placeholder AUTO_SHAPE on master: "
                  f"name='{header_shape.name}', type={header_shape.auto_shape_type} (0.3 pts)")
            total_score += 0.3
        else:
            print("FAIL: Component 1 -- No non-placeholder AUTO_SHAPE found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Shape position and size match expected (0.3 points)
    # Expected: position (0, 0), size 10in x 0.75in
    try:
        if header_shape is not None:
            left_ok = is_approx(header_shape.left, EXPECTED_LEFT)
            top_ok = is_approx(header_shape.top, EXPECTED_TOP)
            width_ok = is_approx(header_shape.width, EXPECTED_WIDTH)
            height_ok = is_approx(header_shape.height, EXPECTED_HEIGHT)

            print(f"  Position: left={header_shape.left} (expected {EXPECTED_LEFT}, ok={left_ok}), "
                  f"top={header_shape.top} (expected {EXPECTED_TOP}, ok={top_ok})")
            print(f"  Size: width={header_shape.width} (expected {EXPECTED_WIDTH}, ok={width_ok}), "
                  f"height={header_shape.height} (expected {EXPECTED_HEIGHT}, ok={height_ok})")

            if left_ok and top_ok and width_ok and height_ok:
                print(f"PASS: Component 2 -- Position and size correct (0.3 pts)")
                total_score += 0.3
            else:
                # Partial: position right but size off, or vice versa
                pos_ok = left_ok and top_ok
                size_ok = width_ok and height_ok
                if pos_ok and not size_ok:
                    print(f"PARTIAL: Component 2 -- Position correct but size incorrect (0.15 pts)")
                    total_score += 0.15
                elif size_ok and not pos_ok:
                    print(f"PARTIAL: Component 2 -- Size correct but position incorrect (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 2 -- Neither position nor size match expected values")
        else:
            print("FAIL: Component 2 -- No header shape found to check position/size")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Shape has solid fill #003366 (0.4 points)
    try:
        if header_shape is not None:
            fill = header_shape.fill
            if fill.type == 1:  # SOLID
                actual_color = str(fill.fore_color.rgb).upper()
                expected_color = EXPECTED_COLOR.upper()
                if actual_color == expected_color:
                    print(f"PASS: Component 3 -- Solid fill color #{actual_color} matches expected #{expected_color} (0.4 pts)")
                    total_score += 0.4
                else:
                    print(f"FAIL: Component 3 -- Fill color #{actual_color} does not match expected #{expected_color}")
            else:
                print(f"FAIL: Component 3 -- Fill type is {fill.type}, expected SOLID (1)")
        else:
            print("FAIL: Component 3 -- No header shape found to check fill color")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
