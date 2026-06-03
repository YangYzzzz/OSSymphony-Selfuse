"""
Initial Setup: Add a consistent header bar to the master slide
Task ID: impress_fix_082
Domain: libreoffice_impress

Creates a 12-slide branded presentation WITHOUT a header bar on the master slide.
The agent's task is to add the dark blue header bar rectangle to the master slide.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_082'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Branded_Deck"
    slide1.placeholders[1].text = "Q1 2025 Strategic Overview\nGlobal Marketing Division"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Analysis & Competitive Landscape"
    items = [
        "Revenue Performance by Region",
        "Customer Acquisition Metrics",
        "Product Roadmap Updates",
        "Team Expansion Plans",
        "Budget Allocation for Q2",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Addressable Market: $4.2B (up 18% YoY)"
    p = body3.add_paragraph()
    p.text = "Key growth drivers: AI adoption, cloud migration, enterprise SaaS expansion"
    p = body3.add_paragraph()
    p.text = "Primary competitors: Acme Corp, NovaTech, ZenithSoft"

    # --- Slide 4: Revenue Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    slide4_title.text_frame.paragraphs[0].text = "Revenue Performance"
    slide4_title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    slide4_title.text_frame.paragraphs[0].runs[0].font.bold = True

    # Add a simple table for revenue data
    tbl_shape = slide4.shapes.add_table(5, 4, Inches(0.5), Inches(1.8), Inches(9), Inches(3))
    tbl = tbl_shape.table
    headers = ["Region", "Q1 Revenue", "Q1 Target", "Variance"]
    for i, h in enumerate(headers):
        tbl.cell(0, i).text = h
    data = [
        ["North America", "$1,245,000", "$1,200,000", "+3.8%"],
        ["Europe", "$892,000", "$950,000", "-6.1%"],
        ["Asia Pacific", "$1,103,000", "$1,000,000", "+10.3%"],
        ["Latin America", "$367,000", "$400,000", "-8.3%"],
    ]
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # --- Slide 5: Customer Acquisition ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Customer Acquisition"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "New enterprise accounts: 47 (target: 40)"
    for line in [
        "Average deal size: $28,500 (up 12% from Q4)",
        "Sales cycle: 62 days average (down from 78 days)",
        "Win rate: 34% (industry avg: 27%)",
        "Top verticals: Healthcare, FinTech, Manufacturing",
    ]:
        p = body5.add_paragraph()
        p.text = line

    # --- Slide 6: Product Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Product Roadmap"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Phase 1 (Complete): Core platform stabilization"
    for line in [
        "Phase 2 (In Progress): AI-powered analytics dashboard",
        "Phase 3 (Q3): Mobile-first redesign",
        "Phase 4 (Q4): Enterprise SSO and compliance suite",
    ]:
        p = body6.add_paragraph()
        p.text = line

    # --- Slide 7: Team Structure ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title7.text_frame.paragraphs[0].text = "Team Structure"
    title7.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title7.text_frame.paragraphs[0].runs[0].font.bold = True

    teams_info = slide7.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf7 = teams_info.text_frame
    tf7.word_wrap = True
    tf7.paragraphs[0].text = "Engineering: 42 engineers (hiring 8 more in Q2)"
    for line in [
        "Product: 12 PMs across 3 product lines",
        "Design: 6 designers, 2 UX researchers",
        "Sales: 28 AEs, 14 SDRs, 4 SEs",
        "Marketing: 15 specialists, 3 content creators",
        "Customer Success: 18 CSMs serving 340 accounts",
    ]:
        p = tf7.add_paragraph()
        p.text = line

    # --- Slide 8: Financial Summary ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Financial Summary"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Total Q1 Revenue: $3,607,000"
    for line in [
        "Operating Expenses: $2,890,000",
        "EBITDA: $717,000 (19.9% margin)",
        "Cash Position: $12.4M",
        "Burn Rate: $480K/month",
        "Runway: 26 months at current burn",
    ]:
        p = body8.add_paragraph()
        p.text = line

    # --- Slide 9: Key Metrics Dashboard ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title9.text_frame.paragraphs[0].text = "Key Metrics Dashboard"
    title9.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title9.text_frame.paragraphs[0].runs[0].font.bold = True

    metrics = slide9.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
    tf9 = metrics.text_frame
    tf9.word_wrap = True
    tf9.paragraphs[0].text = "NPS Score: 72 (up from 65)"
    for line in [
        "Monthly Active Users: 18,400",
        "Churn Rate: 2.1% (target < 3%)",
        "Support Ticket Resolution: 4.2 hours avg",
        "Uptime: 99.97%",
    ]:
        p = tf9.add_paragraph()
        p.text = line

    # --- Slide 10: Competitive Analysis ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Competitive Analysis"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Acme Corp: Strong in enterprise, weak in SMB segment"
    for line in [
        "NovaTech: Aggressive pricing, limited feature set",
        "ZenithSoft: Best UX, struggling with scale",
        "Our advantage: Full-stack solution with superior AI capabilities",
    ]:
        p = body10.add_paragraph()
        p.text = line

    # --- Slide 11: Q2 Priorities ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Q2 Priorities"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "1. Launch AI analytics dashboard (target: June 15)"
    for line in [
        "2. Expand APAC sales team by 6 reps",
        "3. Achieve SOC 2 Type II certification",
        "4. Reduce customer onboarding time to < 5 days",
        "5. Close 3 strategic partnership deals",
    ]:
        p = body11.add_paragraph()
        p.text = line

    # --- Slide 12: Thank You ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[0])
    slide12.shapes.title.text = "Thank You"
    slide12.placeholders[1].text = "Questions & Discussion\nContact: strategy@brandeddeck.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
