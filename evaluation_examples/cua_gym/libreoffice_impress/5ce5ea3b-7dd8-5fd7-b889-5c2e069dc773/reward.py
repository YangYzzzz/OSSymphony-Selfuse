"""
Reward Script: Resize the image on slide 3 to fill the entire slide maintaining its aspect ratio, and center it on the slide.
Task ID: osworld_impress_image_fill_slide_006
Domain: libreoffice_impress
Scoring:
  Component 1: Image fills the slide (covers entire width and height) — 0.4 points
  Component 2: Image is centered on the slide — 0.3 points
  Component 3: Aspect ratio maintained (image placed ratio matches native image aspect ratio) — 0.3 points
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_006'

SLIDE_INDEX = 2  # Slide 3 (0-indexed)
TOLERANCE = 0.005  # 0.5% relative tolerance for position/size comparisons


def is_approx_equal(val1, val2, tolerance=TOLERANCE):
    """Compare two values with a relative tolerance."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def get_image_native_ratio(shape):
    """Get the native (pixel) aspect ratio of an image shape."""
    try:
        from PIL import Image
        import io
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob))
        w, h = img.size
        return w / h if h != 0 else None
    except Exception as e:
        print(f"  WARN: Could not load native image dimensions: {e}")
        return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task: Resize the image on slide 3 to fill the entire slide maintaining its
          aspect ratio, and center it on the slide.
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify 5-slide deck and locate the image on slide 3
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Find the picture shape on slide 3
    picture_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            picture_shape = shape
            break

    if picture_shape is None:
        print("CRITICAL: No image found on slide 3 — cannot score task.")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Slide 3 image found: name='{picture_shape.name}'")
    print(f"INFO: Slide dimensions: {slide_w} x {slide_h} EMU ({slide_w/914400:.2f} x {slide_h/914400:.2f} in)")
    print(f"INFO: Image position: left={picture_shape.left}, top={picture_shape.top}")
    print(f"INFO: Image size: width={picture_shape.width}, height={picture_shape.height}")

    img_left = picture_shape.left
    img_top = picture_shape.top
    img_w = picture_shape.width
    img_h = picture_shape.height

    # -------------------------------------------------------------------------
    # Component 1: Image fills (covers) the entire slide (0.4 points)
    #
    # The image must cover the full slide dimensions. The image width and height
    # must each be >= the slide dimensions (within tolerance), to ensure it fills.
    # -------------------------------------------------------------------------
    try:
        fills_width = img_w >= slide_w or is_approx_equal(img_w, slide_w)
        fills_height = img_h >= slide_h or is_approx_equal(img_h, slide_h)

        if fills_width and fills_height:
            print(f"PASS: Component 1 — Image fills the entire slide: "
                  f"img_w={img_w/914400:.4f}in >= slide_w={slide_w/914400:.4f}in, "
                  f"img_h={img_h/914400:.4f}in >= slide_h={slide_h/914400:.4f}in (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Image does not fill the slide: "
                  f"img_w={img_w/914400:.4f}in, slide_w={slide_w/914400:.4f}in, "
                  f"img_h={img_h/914400:.4f}in, slide_h={slide_h/914400:.4f}in")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------------
    # Component 2: Image is centered on the slide (0.3 points)
    #
    # The expected centered position is:
    #   expected_left = (slide_w - img_w) / 2
    #   expected_top  = (slide_h - img_h) / 2
    # When the image exactly fills the slide (img_w == slide_w, img_h == slide_h),
    # the centered position is (0, 0).
    # -------------------------------------------------------------------------
    try:
        expected_left = (slide_w - img_w) // 2
        expected_top = (slide_h - img_h) // 2

        left_ok = is_approx_equal(img_left, expected_left) or (expected_left == 0 and img_left == 0)
        top_ok = is_approx_equal(img_top, expected_top) or (expected_top == 0 and img_top == 0)

        if left_ok and top_ok:
            print(f"PASS: Component 2 — Image is centered: "
                  f"left={img_left} (expected {expected_left}), "
                  f"top={img_top} (expected {expected_top}) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Image is not centered: "
                  f"left={img_left} (expected ~{expected_left}), "
                  f"top={img_top} (expected ~{expected_top})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------------------------
    # Component 3: Aspect ratio maintained AND image fills the entire slide (0.3 points)
    #
    # The image must fill the full slide AND the placed aspect ratio must match
    # the native image pixel ratio within tolerance (0.5%). This compound check
    # ensures the resize was proportional while covering the slide.
    # An image that merely maintains aspect ratio in its initial smaller size
    # will NOT pass this component — it must both fill AND maintain aspect ratio.
    # -------------------------------------------------------------------------
    try:
        if img_h == 0:
            print("FAIL: Component 3 — Image height is 0, cannot compute aspect ratio")
        else:
            placed_ratio = img_w / img_h
            native_ratio = get_image_native_ratio(picture_shape)

            # Both conditions must hold: fills the slide AND ratio is maintained
            fills_check = (img_w >= slide_w or is_approx_equal(img_w, slide_w)) and \
                          (img_h >= slide_h or is_approx_equal(img_h, slide_h))

            if native_ratio is None:
                ratio_check = is_approx_equal(placed_ratio, slide_w / slide_h)
            else:
                ratio_check = is_approx_equal(placed_ratio, native_ratio)

            if fills_check and ratio_check:
                print(f"PASS: Component 3 — Aspect ratio maintained AND image fills slide: "
                      f"placed ratio {placed_ratio:.4f} matches native ratio "
                      f"{native_ratio if native_ratio else slide_w/slide_h:.4f}, "
                      f"img fills slide (0.3 pts)")
                total_score += 0.3
            elif not fills_check:
                print(f"FAIL: Component 3 — Image does not fill slide (required for this component): "
                      f"img_w={img_w/914400:.4f}in, slide_w={slide_w/914400:.4f}in, "
                      f"img_h={img_h/914400:.4f}in, slide_h={slide_h/914400:.4f}in")
            else:
                print(f"FAIL: Component 3 — Aspect ratio not maintained while filling: "
                      f"placed ratio {placed_ratio:.4f} != native ratio "
                      f"{native_ratio if native_ratio else slide_w/slide_h:.4f}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
