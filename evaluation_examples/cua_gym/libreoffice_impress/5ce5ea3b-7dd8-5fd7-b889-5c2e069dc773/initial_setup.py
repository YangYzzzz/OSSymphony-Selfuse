"""
Initial Setup: 5-slide architecture portfolio deck with building photograph on slide 3
Task ID: osworld_impress_image_fill_slide_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_fill_slide_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMAGE_PATH = f'{WORKDIR}/{TASK_ID}_building.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_building_image(path: str, width: int = 800, height: int = 600):
    """Create a realistic-looking building photograph placeholder."""
    img = Image.new('RGB', (width, height), color=(135, 160, 185))
    draw = ImageDraw.Draw(img)

    # Sky gradient effect
    for y in range(height // 3):
        r = int(100 + y * 0.3)
        g = int(140 + y * 0.2)
        b = int(200 - y * 0.1)
        draw.line([(0, y), (width, y)], fill=(r, g, b))

    # Ground
    draw.rectangle([0, height * 2 // 3, width, height], fill=(80, 100, 70))

    # Main building facade
    bldg_left = width // 6
    bldg_right = width * 5 // 6
    bldg_top = height // 8
    bldg_bottom = height * 2 // 3
    draw.rectangle([bldg_left, bldg_top, bldg_right, bldg_bottom], fill=(220, 215, 200))
    draw.rectangle([bldg_left, bldg_top, bldg_right, bldg_bottom], outline=(160, 150, 130), width=3)

    # Windows grid
    win_cols = 8
    win_rows = 6
    win_w = (bldg_right - bldg_left - 40) // win_cols
    win_h = (bldg_bottom - bldg_top - 40) // win_rows
    for row in range(win_rows):
        for col in range(win_cols):
            wx = bldg_left + 20 + col * win_w
            wy = bldg_top + 30 + row * win_h
            draw.rectangle([wx, wy, wx + win_w - 8, wy + win_h - 8], fill=(160, 190, 210), outline=(100, 120, 140), width=1)

    # Entrance
    ent_w = 60
    ent_h = 90
    ent_x = (bldg_left + bldg_right) // 2 - ent_w // 2
    ent_y = bldg_bottom - ent_h
    draw.rectangle([ent_x, ent_y, ent_x + ent_w, bldg_bottom], fill=(60, 50, 40))
    draw.rectangle([ent_x, ent_y, ent_x + ent_w, ent_y + 10], fill=(100, 80, 60))

    # Roofline details
    draw.rectangle([bldg_left - 5, bldg_top - 10, bldg_right + 5, bldg_top], fill=(180, 170, 150))

    # Trees / landscaping
    for tx in [bldg_left - 40, bldg_right + 10]:
        draw.ellipse([tx - 20, height // 2, tx + 40, height * 2 // 3 + 10], fill=(50, 120, 60))
        draw.rectangle([tx + 5, height * 2 // 3, tx + 15, height * 2 // 3 + 20], fill=(80, 60, 40))

    # Label
    draw.text((10, height - 25), "Meridian Tower — Architecture Portfolio", fill=(240, 240, 240))

    img.save(path, 'PNG')
    print(f'Building image created: {path}')


def create_initial():
    # Create the building image first
    create_building_image(IMAGE_PATH, width=800, height=600)

    prs = Presentation()
    # Standard 10x7.5 inch widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Meridian Architecture Portfolio"
    slide1.placeholders[1].text = "Urban Design & Structural Innovation 2025"
    # Background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3A)
    # Title formatting
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xF0, 0xE8, 0xD8)
        run.font.size = Pt(36)
        run.font.bold = True
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE0)
        run.font.size = Pt(20)

    # ── Slide 2: Practice Overview ────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Our Practice"
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF0)
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Founded in 1998, Meridian Architects has delivered award-winning projects across 30+ countries."
    lines2 = [
        "Commercial & Mixed-Use Developments",
        "Sustainable Green Building Design",
        "Urban Masterplanning",
        "Heritage Conservation & Adaptive Reuse",
    ]
    for line in lines2:
        p = tf2.add_paragraph()
        p.text = line
        p.level = 1
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1A, 0x2A, 0x3A)
            run.font.size = Pt(32)
            run.font.bold = True

    # ── Slide 3: Featured Project — Building Photograph ───────────────────────
    # Image is SMALLER than the slide (not filling) — this is the initial state
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x10, 0x18, 0x24)

    # Title text box
    txBox3 = slide3.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.7))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Meridian Tower — Featured Project"
    p3.alignment = PP_ALIGN.CENTER
    for run in p3.runs:
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xF0, 0xE8, 0xD8)

    # Building image — smaller than slide, placed in the middle-ish area
    # Image native: 800x600 px → aspect ratio 4:3
    # Place at 3" x 2.5" with size 4" x 3" (not filling the 10x7.5 slide)
    img_w = Inches(4)
    img_h = Inches(3)
    img_left = (slide_w - img_w) // 2   # centered horizontally
    img_top = Inches(1.5)               # below title, not filling
    pic3 = slide3.shapes.add_picture(IMAGE_PATH, img_left, img_top, img_w, img_h)

    # Caption text box below image
    cap_top = img_top + img_h + Inches(0.15)
    txBox3c = slide3.shapes.add_textbox(Inches(1.0), cap_top, Inches(8.0), Inches(0.5))
    tfc = txBox3c.text_frame
    pc = tfc.paragraphs[0]
    pc.text = "42-story mixed-use tower, Cape Town | Completed 2023"
    pc.alignment = PP_ALIGN.CENTER
    for run in pc.runs:
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE0)

    # ── Slide 4: Design Philosophy ────────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Design Philosophy"
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF0)
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "We believe architecture should harmonize human experience with the environment."
    philosophy_points = [
        "Biophilic design principles embedded in every project",
        "Net-zero carbon targets by 2030",
        "Community-centered spatial planning",
        "Material traceability & circular economy alignment",
    ]
    for pt in philosophy_points:
        p = tf4.add_paragraph()
        p.text = pt
        p.level = 1
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0x1A, 0x2A, 0x3A)
            run.font.size = Pt(32)
            run.font.bold = True

    # ── Slide 5: Contact ──────────────────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[0])
    slide5.shapes.title.text = "Get In Touch"
    slide5.placeholders[1].text = "studio@meridianarch.com\n+27 21 555 0192\nwww.meridianarch.com"
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x3A)
    for run in slide5.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xF0, 0xE8, 0xD8)
        run.font.size = Pt(36)
        run.font.bold = True
    for para in slide5.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xB0, 0xC8, 0xE0)
            run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
