"""
Reward Script: Insert a blue rectangle with 'Call to Action' white text on slide 3
Task ID: impress_tm_048
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Rectangle/auto shape exists on slide 3 (beyond the original 2 shapes)
  Component 2 (0.3): Shape contains text 'Call to Action'
  Component 3 (0.2): Shape fill is solid blue (#0066CC)
  Component 4 (0.2): Text color is white (#FFFFFF)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_048'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3

    # Find any non-placeholder auto shapes on slide 3 (the task-added rectangle)
    # Original slide 3 has only PLACEHOLDER shapes; the task adds an AUTO_SHAPE
    candidate_shapes = []
    for shape in slide3.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            candidate_shapes.append(shape)

    # Component 1: Rectangle/auto shape exists on slide 3 (0.3 points)
    try:
        if len(candidate_shapes) > 0:
            print(f"PASS: Component 1 - Found {len(candidate_shapes)} auto shape(s) on slide 3 (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - No auto shapes found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if len(candidate_shapes) == 0:
        # No shape to check further components against
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Find the best candidate shape (one that contains 'Call to Action' text, or first one)
    target_shape = None
    for shape in candidate_shapes:
        if shape.has_text_frame:
            full_text = ""
            for para in shape.text_frame.paragraphs:
                full_text += para.text
            if "call to action" in full_text.lower():
                target_shape = shape
                break
    if target_shape is None:
        target_shape = candidate_shapes[0]

    # Component 2: Shape contains text 'Call to Action' (0.3 points)
    try:
        if target_shape.has_text_frame:
            full_text = ""
            for para in target_shape.text_frame.paragraphs:
                full_text += para.text
            if "call to action" in full_text.lower():
                print(f"PASS: Component 2 - Shape text contains 'Call to Action': {repr(full_text)} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 - Shape text is {repr(full_text)}, expected 'Call to Action'")
        else:
            print(f"FAIL: Component 2 - Target shape has no text frame")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Shape fill is solid blue #0066CC (0.2 points)
    try:
        fill = target_shape.fill
        if fill.type == 1:  # solid fill
            fill_color = str(fill.fore_color.rgb).upper()
            if fill_color == "0066CC":
                print(f"PASS: Component 3 - Shape fill is solid blue #0066CC (0.2 pts)")
                total_score += 0.2
            else:
                # Accept close blue variants
                print(f"FAIL: Component 3 - Shape fill color is #{fill_color}, expected #0066CC")
        else:
            print(f"FAIL: Component 3 - Shape fill type is {fill.type}, expected solid (1)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Text color is white #FFFFFF (0.2 points)
    try:
        white_found = False
        if target_shape.has_text_frame:
            for para in target_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        try:
                            if run.font.color.type is not None:
                                color_rgb = str(run.font.color.rgb).upper()
                                if color_rgb == "FFFFFF":
                                    white_found = True
                        except Exception:
                            pass
        if white_found:
            print(f"PASS: Component 4 - Text color is white #FFFFFF (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 - Text color is not white #FFFFFF")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
