"""
Initial Setup: Insert a rectangular shape on slide 3 with 'Call to Action' text
Task ID: impress_tm_048
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_048'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Marketing Strategy"
    slide1.placeholders[1].text = "Digital Campaigns & Brand Positioning"

    # --- Slide 2: Campaign Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Campaign Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Social media engagement increased by 34% in Q2"
    p2a = tf2.add_paragraph()
    p2a.text = "Email open rates averaged 22.5% across all segments"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "Influencer partnerships drove 15,000 new sign-ups"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "Paid search ROI improved to 3.8x from 2.9x"
    p2c.level = 0

    # --- Slide 3: Next Steps (NO rectangle, NO 'Call to Action') ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Next Steps"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Launch retargeting campaign by August 15th"
    p3a = tf3.add_paragraph()
    p3a.text = "Finalize creative assets for fall product line"
    p3a.level = 0
    p3b = tf3.add_paragraph()
    p3b.text = "Schedule stakeholder review meetings for Week 32"
    p3b.level = 0
    p3c = tf3.add_paragraph()
    p3c.text = "Allocate remaining Q3 budget across top channels"
    p3c.level = 0
    p3d = tf3.add_paragraph()
    p3d.text = "Prepare A/B testing plan for landing page variants"
    p3d.level = 0

    # --- Slide 4: Budget Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Budget Breakdown"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Total Q3 Budget: $285,000"
    p4a = tf4.add_paragraph()
    p4a.text = "Social Media Ads: $95,000 (33%)"
    p4a.level = 0
    p4b = tf4.add_paragraph()
    p4b.text = "Search & Display: $72,000 (25%)"
    p4b.level = 0
    p4c = tf4.add_paragraph()
    p4c.text = "Content Production: $58,000 (20%)"
    p4c.level = 0
    p4d = tf4.add_paragraph()
    p4d.text = "Influencer Program: $38,000 (13%)"
    p4d.level = 0
    p4e = tf4.add_paragraph()
    p4e.text = "Events & Sponsorships: $22,000 (8%)"
    p4e.level = 0

    # --- Slide 5: Timeline ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Timeline"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "July 1-15: Creative brief finalization"
    p5a = tf5.add_paragraph()
    p5a.text = "July 16-31: Asset production and review"
    p5a.level = 0
    p5b = tf5.add_paragraph()
    p5b.text = "August 1-14: Campaign setup and QA"
    p5b.level = 0
    p5c = tf5.add_paragraph()
    p5c.text = "August 15: Full campaign launch"
    p5c.level = 0
    p5d = tf5.add_paragraph()
    p5d.text = "September 30: Q3 performance review"
    p5d.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
