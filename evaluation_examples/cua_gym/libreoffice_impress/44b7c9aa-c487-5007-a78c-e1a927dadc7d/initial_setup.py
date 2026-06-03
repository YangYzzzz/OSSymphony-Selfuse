"""
Initial Setup: Format notes on slide 2 of a strategy presentation
Task ID: impress_ndo_012
Domain: libreoffice_impress

Creates a multi-slide strategy presentation. Slide 2 has speaker notes
in default formatting (12pt, regular, black). The agent must reformat them.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Strategic Review"
    slide1.placeholders[1].text = "Global Operations Division\nPrepared by Sarah Chen, VP Strategy"

    # ---- Slide 2: Strategy Overview (the target slide) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Customer Retention Strategy"

    # Content body
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.clear()

    items = [
        "Current retention rate: 72% (industry avg: 68%)",
        "Target: Achieve 30% improvement by end of Q4",
        "Key driver: Personalized engagement workflows",
        "Budget allocation: $1.2M across three initiatives",
        "Timeline: Phase 1 rollout begins June 2025",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    # Speaker notes for slide 2 -- default formatting (12pt, regular, black)
    notes_slide = slide2.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.clear()

    note_lines = [
        "Key Strategy Points",
        "Focus on customer retention rates. Emphasize the 30% improvement target. Reference the competitive analysis from Q1.",
    ]

    for i, line in enumerate(note_lines):
        if i == 0:
            p = notes_tf.paragraphs[0]
        else:
            p = notes_tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.italic = False
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ---- Slide 3: Competitive Landscape ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Competitive Landscape Analysis"

    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()

    comp_items = [
        "Competitor A: Launched loyalty program (est. 15% lift)",
        "Competitor B: Acquired CRM startup for $400M",
        "Market shift toward subscription-based models",
        "Our advantage: Proprietary data pipeline from 2.3M users",
    ]
    for i, item in enumerate(comp_items):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    # ---- Slide 4: Financial Projections ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Projections"

    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.clear()

    fin_items = [
        "Projected revenue uplift: $3.8M annually",
        "Cost of retention programs: $1.2M (ROI: 3.2x)",
        "Break-even expected by Q3 2025",
        "Net customer lifetime value increase: 22%",
    ]
    for i, item in enumerate(fin_items):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
