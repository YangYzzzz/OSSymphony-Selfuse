"""
Reward Script: Format notes on slide 2 of Strategy.pptx
Task ID: impress_ndo_012
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): First line bold and 16pt
  Component 2 (0.35): Remaining text italic and 11pt
  Component 3 (0.15): Remaining text color #333333
  Component 4 (0.15): Text content preserved unchanged
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_012'

EXPECTED_FIRST_LINE = 'Key Strategy Points'
EXPECTED_REMAINING = 'Focus on customer retention rates. Emphasize the 30% improvement target. Reference the competitive analysis from Q1.'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)

    # Precondition: slide must have notes
    try:
        ns = slide.notes_slide
        tf = ns.notes_text_frame
        paragraphs = tf.paragraphs
    except Exception as e:
        print(f"FAIL: Cannot access notes on slide 2: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(paragraphs) < 2:
        print(f"FAIL: Notes have {len(paragraphs)} paragraphs, expected at least 2")
        print("REWARD: 0.0")
        return 0.0

    first_para = paragraphs[0]
    remaining_para = paragraphs[1]

    # Component 1: First line bold and 16pt (0.35 points)
    # Initial: bold=False, size=152400 (12pt). Golden: bold=True, size=203200 (16pt).
    try:
        runs = [r for r in first_para.runs if (r.text or "").strip()]
        if len(runs) == 0:
            print("FAIL: Component 1 -- first paragraph has no non-empty runs")
        else:
            all_bold = True
            all_16pt = True
            for run in runs:
                b = run.font.bold
                if b is None or b is False:
                    all_bold = False
                sz = run.font.size
                # 16pt = 203200 EMU
                if sz is None or abs(sz - 203200) > 1000:
                    all_16pt = False

            if all_bold and all_16pt:
                print(f"PASS: Component 1 -- first line is bold and 16pt (0.35 pts)")
                total_score += 0.35
            elif all_bold:
                print(f"PARTIAL: Component 1 -- first line is bold but not 16pt (0.15 pts)")
                total_score += 0.15
            elif all_16pt:
                print(f"PARTIAL: Component 1 -- first line is 16pt but not bold (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 -- first line bold={all_bold}, 16pt={all_16pt}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Remaining text italic and 11pt (0.35 points)
    # Initial: italic=False, size=152400 (12pt). Golden: italic=True, size=139700 (11pt).
    try:
        runs = [r for r in remaining_para.runs if (r.text or "").strip()]
        if len(runs) == 0:
            print("FAIL: Component 2 -- remaining paragraph has no non-empty runs")
        else:
            all_italic = True
            all_11pt = True
            for run in runs:
                it = run.font.italic
                if it is None or it is False:
                    all_italic = False
                sz = run.font.size
                # 11pt = 139700 EMU
                if sz is None or abs(sz - 139700) > 1000:
                    all_11pt = False

            if all_italic and all_11pt:
                print(f"PASS: Component 2 -- remaining text is italic and 11pt (0.35 pts)")
                total_score += 0.35
            elif all_italic:
                print(f"PARTIAL: Component 2 -- remaining text is italic but not 11pt (0.15 pts)")
                total_score += 0.15
            elif all_11pt:
                print(f"PARTIAL: Component 2 -- remaining text is 11pt but not italic (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 -- remaining text italic={all_italic}, 11pt={all_11pt}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Remaining text color #333333 (0.15 points)
    # Initial: color=000000. Golden: color=333333.
    try:
        runs = [r for r in remaining_para.runs if (r.text or "").strip()]
        if len(runs) == 0:
            print("FAIL: Component 3 -- remaining paragraph has no non-empty runs")
        else:
            all_correct_color = True
            for run in runs:
                try:
                    if run.font.color.type is not None:
                        rgb_str = str(run.font.color.rgb).upper()
                        if rgb_str != '333333':
                            all_correct_color = False
                            print(f"  DEBUG: run color={rgb_str}, expected 333333")
                    else:
                        all_correct_color = False
                except Exception:
                    all_correct_color = False

            if all_correct_color:
                print(f"PASS: Component 3 -- remaining text color is #333333 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- remaining text color is not #333333")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Text content preserved unchanged (0.15 points)
    # Both initial and golden should have the same text content.
    # This component verifies that formatting changes did NOT alter text.
    # It scores 0 on initial because it's combined with a task-change check:
    # text unchanged AND at least one formatting change detected.
    try:
        first_text = first_para.text.strip()
        remaining_text = remaining_para.text.strip()

        text_ok = (first_text == EXPECTED_FIRST_LINE and remaining_text == EXPECTED_REMAINING)

        # Gate: at least one formatting change must be present (so initial scores 0)
        any_format_change = False
        for run in first_para.runs:
            if run.font.bold is True:
                any_format_change = True
                break
        if not any_format_change:
            for run in remaining_para.runs:
                if run.font.italic is True:
                    any_format_change = True
                    break

        if text_ok and any_format_change:
            print(f"PASS: Component 4 -- text content preserved and formatting applied (0.15 pts)")
            total_score += 0.15
        elif text_ok and not any_format_change:
            print(f"FAIL: Component 4 -- text correct but no formatting changes detected (initial state)")
        else:
            print(f"FAIL: Component 4 -- text content changed. First='{first_text}', Remaining='{remaining_text}'")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state("libreoffice_impress")
    verify_task(file_path)
