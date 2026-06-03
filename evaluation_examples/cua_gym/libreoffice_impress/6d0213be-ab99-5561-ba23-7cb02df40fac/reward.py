"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 142 I need Picture 1 to sit flush with the content column—specifically, align it to the left margin so the picture’s left edge is exactly 1.0 cm in from the slide border. How do I do that in LibreOffice Impress?
Generated: 2025-09-10 15:52:42
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation


def verify_picture_alignment(
    file_path: str,
    slide_number: int = 142,
    expected_margin_cm: float = 1.0,
    tolerance_cm_exact: float = 0.05,
    tolerance_cm_loose: float = 0.5,
):
    """Reward-script verification for the LibreOffice-Impress task.

    Task recap:
        On slide 142, Picture 1 must be aligned so its **left edge** is exactly
        1.0 cm from the slide border (i.e. flush with the content column).

    Scoring (progressive):
        • 0.4 pts – Picture 1 is present on slide 142.
        • 0.6 pts – Alignment accuracy:
              – +0.6 if within ±0.05 cm (exact tolerance)
              – +0.3 if within ±0.5 cm (loose tolerance)
              – 0   otherwise
    The final score is capped at 1.0.
    """

    print(f"Starting verification for file: {file_path}")
    max_score = 1.0
    score = 0.0

    # Constants – EMU (English Metric Unit) conversion (1 cm = 360 000 EMU)
    EMU_PER_CM = 360_000
    expected_left = int(expected_margin_cm * EMU_PER_CM)
    tol_exact = int(tolerance_cm_exact * EMU_PER_CM)
    tol_loose = int(tolerance_cm_loose * EMU_PER_CM)

    # ------------------------------------------------------------------
    # 1. Load presentation (no points for merely existing / loading)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – task failed")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded – {len(prs.slides)} slides detected")
    except Exception as exc:
        print(f"✗ Error loading presentation: {exc}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure slide 142 exists
    # ------------------------------------------------------------------
    if len(prs.slides) < slide_number:
        print(
            f"✗ Slide {slide_number} missing (only {len(prs.slides)} slides present)"
        )
        return 0.0

    slide = prs.slides[slide_number - 1]  # zero-based index

    # ------------------------------------------------------------------
    # 3. Locate “Picture 1” on that slide
    # ------------------------------------------------------------------
    picture = None
    for shp in slide.shapes:
        if shp.shape_type == 13 and shp.name.strip().lower() == "picture 1":
            picture = shp
            break

    if picture is None:
        print("✗ ‘Picture 1’ not found on slide – no alignment possible")
    else:
        print(
            f"✓ Found ‘Picture 1’ – left position: {picture.left} EMU "
            f"({picture.left / EMU_PER_CM:.2f} cm)"
        )
        score += 0.4  # picture exists

        # ------------------------------------------------------------------
        # 4. Verify alignment accuracy
        # ------------------------------------------------------------------
        diff = abs(picture.left - expected_left)
        print(
            f"Alignment difference: {diff} EMU ({diff / EMU_PER_CM:.4f} cm)"
        )

        if diff <= tol_exact:
            print("✓ Alignment within exact tolerance (±0.05 cm)")
            score += 0.6
        elif diff <= tol_loose:
            print("⚠ Alignment within loose tolerance (±0.5 cm) – partial credit")
            score += 0.3
        else:
            print("✗ Alignment outside acceptable tolerance – no credit")

    final_score = min(score, max_score)
    print(f"Total SCORE: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path to the file inside the VM workspace
    FILE_PATH = "/home/user/on_slide_142_i_need_picture_1_to_sit_flush_with_the_content_columnspecifically_align_it_to_the_left__golden.pptx"

    reward = verify_picture_alignment(FILE_PATH)
    print(f"REWARD: {reward}")

