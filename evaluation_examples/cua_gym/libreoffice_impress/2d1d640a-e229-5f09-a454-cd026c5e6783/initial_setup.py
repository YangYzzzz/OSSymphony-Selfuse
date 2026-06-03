"""
Initial Setup: Set a background image on the master slide using texture_bg.jpg tiled across the entire slide.
Task ID: impress_ma_016
Domain: libreoffice_impress

Creates:
  - /home/user/impress_ma_016.pptx  (7-slide Creative Brief presentation with light gray master bg)
  - /home/user/Desktop/texture_bg.jpg  (300x300 seamless tile pattern)
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw
import numpy as np

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
TEXTURE_PATH = f'{WORKDIR}/Desktop/texture_bg.jpg'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_texture_image():
    """Create a 300x300 seamless tileable texture pattern."""
    os.makedirs(os.path.dirname(TEXTURE_PATH), exist_ok=True)
    size = 300
    img = Image.new('RGB', (size, size), (210, 200, 185))
    draw = ImageDraw.Draw(img)

    # Create a subtle woven/fabric-like texture pattern
    np.random.seed(42)
    for y in range(0, size, 6):
        for x in range(0, size, 6):
            r = 195 + np.random.randint(-15, 16)
            g = 185 + np.random.randint(-15, 16)
            b = 170 + np.random.randint(-15, 16)
            draw.rectangle([x, y, x + 5, y + 5], fill=(r, g, b))

    # Add subtle diagonal lines for texture
    for i in range(-size, size * 2, 12):
        draw.line([(i, 0), (i + size, size)], fill=(200, 190, 175), width=1)

    img.save(TEXTURE_PATH, 'JPEG', quality=90)
    print(f'Texture image created: {TEXTURE_PATH}')


def create_presentation():
    """Create a 7-slide Creative Brief presentation with light gray master background."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Creative Brief: Q2 Brand Refresh"
    slide1.placeholders[1].text = "Prepared by: Aria Nakamura | Design Lead\nApril 2025"

    # --- Slide 2: Project Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Objective: Refresh the Meridian Corp visual identity for the summer product launch."
    p = body2.add_paragraph()
    p.text = "Timeline: April 14 – June 30, 2025"
    p.level = 0
    p = body2.add_paragraph()
    p.text = "Budget: $85,000 allocated across digital, print, and social channels"
    p.level = 0
    p = body2.add_paragraph()
    p.text = "Stakeholders: Marketing VP, Product Design, External Agency (Bloom Studio)"
    p.level = 0

    # --- Slide 3: Target Audience ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Target Audience"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Primary: Professionals aged 28–45 in urban markets"
    p = body3.add_paragraph()
    p.text = "Secondary: Design-conscious consumers seeking premium lifestyle products"
    p.level = 0
    p = body3.add_paragraph()
    p.text = "Key insight: 72% of target demographic engages with brand stories on Instagram and LinkedIn"
    p.level = 0
    p = body3.add_paragraph()
    p.text = "Psychographics: Value authenticity, sustainability, and minimalist aesthetics"
    p.level = 0

    # --- Slide 4: Brand Messaging ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Brand Messaging Framework"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Tagline: \"Crafted for the Thoughtful.\""
    p = body4.add_paragraph()
    p.text = "Tone: Warm, confident, understated elegance"
    p.level = 0
    p = body4.add_paragraph()
    p.text = "Core Values: Precision, Sustainability, Innovation"
    p.level = 0
    p = body4.add_paragraph()
    p.text = "Voice Guidelines: First-person plural (\"we\"), active voice, avoid jargon"
    p.level = 0

    # --- Slide 5: Visual Direction ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Visual Direction"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Color Palette: Warm neutrals (#C8B89A, #E8DFD0) with accent teal (#2E8B8B)"
    p = body5.add_paragraph()
    p.text = "Typography: Montserrat (headings), Source Sans Pro (body)"
    p.level = 0
    p = body5.add_paragraph()
    p.text = "Photography: Natural lighting, earth tones, lifestyle-oriented compositions"
    p.level = 0
    p = body5.add_paragraph()
    p.text = "Texture & Pattern: Organic materials, linen, handmade paper references"
    p.level = 0

    # --- Slide 6: Deliverables ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Deliverables & Timeline"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Phase 1 (Apr 14–30): Mood boards, color studies, 3 logo concepts"
    p = body6.add_paragraph()
    p.text = "Phase 2 (May 1–31): Selected concept refinement, brand guidelines draft"
    p.level = 0
    p = body6.add_paragraph()
    p.text = "Phase 3 (Jun 1–15): Final assets for print and digital, social templates"
    p.level = 0
    p = body6.add_paragraph()
    p.text = "Phase 4 (Jun 16–30): Launch support, A/B testing creative variants"
    p.level = 0

    # --- Slide 7: Next Steps ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Next Steps"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "1. Schedule kickoff meeting with Bloom Studio (week of April 14)"
    p = body7.add_paragraph()
    p.text = "2. Distribute audience research summary to stakeholders"
    p.level = 0
    p = body7.add_paragraph()
    p.text = "3. Set up shared workspace in Figma for collaborative design review"
    p.level = 0
    p = body7.add_paragraph()
    p.text = "4. Confirm budget sign-off from VP Marketing by April 11"
    p.level = 0

    # Save initial version
    prs.save(OUTPUT)

    # Now modify the slide master background to solid light gray via XML
    _set_master_background_gray(OUTPUT)

    print(f'Initial presentation created: {OUTPUT}')


def _set_master_background_gray(pptx_path):
    """Set the slide master background to solid light gray (#D3D3D3) via XML manipulation."""
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }

    # Register namespaces to avoid ns0/ns1 prefixes
    for prefix, uri in ns.items():
        ET.register_namespace(prefix, uri)
    ET.register_namespace('', 'http://schemas.openxmlformats.org/presentationml/2006/main')

    temp_path = pptx_path + '.tmp'
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(temp_path, 'w') as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == 'ppt/slideMasters/slideMaster1.xml':
                    root = ET.fromstring(data)
                    # Find or create cSld/bg
                    cSld = root.find('{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
                    if cSld is None:
                        cSld = ET.SubElement(root, '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')

                    # Remove existing bg if any
                    bg = cSld.find('{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                    if bg is not None:
                        cSld.remove(bg)

                    # Create solid gray background
                    bg = ET.SubElement(cSld, '{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                    bgPr = ET.SubElement(bg, '{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr')
                    solidFill = ET.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                    srgbClr = ET.SubElement(solidFill, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    srgbClr.set('val', 'D3D3D3')
                    ET.SubElement(bgPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst')

                    # bg must be the first child of cSld
                    cSld.remove(bg)
                    cSld.insert(0, bg)

                    data = ET.tostring(root, xml_declaration=True, encoding='UTF-8')
                zout.writestr(item, data)

    os.replace(temp_path, pptx_path)
    print('Master slide background set to light gray (#D3D3D3)')


def main():
    create_texture_image()
    create_presentation()

    # GUI-ready: open the presentation in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
