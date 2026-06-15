"""
Initial Setup: Insert a right-arrow connector shape between two boxes on slide 4.
Task ID: impress_tm_060
Domain: libreoffice_impress

Creates a 6-slide business process presentation. Slide 4 has two rectangular
shapes ('Input' on the left, 'Output' on the right) with empty space between them.
No connector arrow exists yet -- that is the task for the agent.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_060'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, title_text, subtitle_text=None):
    """Helper to set title and optional subtitle on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Try to set subtitle in placeholder 1
    if subtitle_text:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = subtitle_text
                break


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a simple textbox to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    add_title_text(slide1, "Data Processing Pipeline Overview",
                   "Q2 2025 Engineering Review")

    # ===== Slide 2: Agenda =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Agenda")
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "1. Current Architecture"
    items = [
        "2. Data Flow Analysis",
        "3. Input/Output Pipeline",
        "4. Process Flow Diagram",
        "5. Performance Metrics",
        "6. Next Steps & Timeline",
    ]
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0

    # ===== Slide 3: Architecture Overview =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Current Architecture", font_size=28, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))
    add_textbox(slide3, Inches(0.5), Inches(1.5), Inches(12), Inches(1.2),
                "The pipeline processes approximately 2.3 million records daily "
                "through three main stages: ingestion, transformation, and delivery. "
                "Average latency has improved from 450ms to 120ms since the Q1 refactor.",
                font_size=16)
    add_textbox(slide3, Inches(0.5), Inches(3.0), Inches(12), Inches(3),
                "Key Components:\n"
                "  - Apache Kafka message broker (v3.6)\n"
                "  - Spark Streaming cluster (24 nodes)\n"
                "  - PostgreSQL 16 data warehouse\n"
                "  - Redis cache layer (128 GB)\n"
                "  - Grafana monitoring dashboard",
                font_size=14)

    # ===== Slide 4: Process Flow - Input and Output boxes (KEY SLIDE) =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Process Flow Diagram", font_size=28, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))

    # "Input" box on the left side
    input_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(3.0),   # left, top
        Inches(3.0), Inches(2.0),   # width, height
    )
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(0x4E, 0x79, 0xA7)
    input_box.line.color.rgb = RGBColor(0x2C, 0x5F, 0x8A)
    input_box.line.width = Pt(2)
    tf_in = input_box.text_frame
    tf_in.word_wrap = True
    p_in = tf_in.paragraphs[0]
    p_in.text = "Input"
    p_in.alignment = PP_ALIGN.CENTER
    run_in = p_in.runs[0]
    run_in.font.size = Pt(24)
    run_in.font.bold = True
    run_in.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # "Output" box on the right side
    output_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(3.0),   # left, top
        Inches(3.0), Inches(2.0),   # width, height
    )
    output_box.fill.solid()
    output_box.fill.fore_color.rgb = RGBColor(0x59, 0xA1, 0x4F)
    output_box.line.color.rgb = RGBColor(0x3D, 0x7A, 0x36)
    output_box.line.width = Pt(2)
    tf_out = output_box.text_frame
    tf_out.word_wrap = True
    p_out = tf_out.paragraphs[0]
    p_out.text = "Output"
    p_out.alignment = PP_ALIGN.CENTER
    run_out = p_out.runs[0]
    run_out.font.size = Pt(24)
    run_out.font.bold = True
    run_out.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add description text below the boxes
    add_textbox(slide4, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
                "The process flow requires a connector arrow between the Input and Output "
                "stages to visually represent the data transformation pipeline.",
                font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ===== Slide 5: Performance Metrics =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(10), Inches(1),
                "Performance Metrics", font_size=28, bold=True,
                color=RGBColor(0x1F, 0x49, 0x7D))

    # Add a table with metrics
    table_shape = slide5.shapes.add_table(
        5, 4, Inches(1), Inches(1.5), Inches(10), Inches(3))
    table = table_shape.table
    headers = ["Metric", "Q1 2025", "Q2 2025", "Change"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    metrics_data = [
        ["Throughput (rec/sec)", "18,400", "26,750", "+45.4%"],
        ["Avg Latency (ms)", "450", "120", "-73.3%"],
        ["Error Rate (%)", "0.34", "0.08", "-76.5%"],
        ["Uptime (%)", "99.82", "99.97", "+0.15%"],
    ]
    for r, row_data in enumerate(metrics_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # ===== Slide 6: Next Steps =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide6, "Next Steps & Timeline")
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Migrate remaining batch jobs to streaming (June 2025)"
    next_items = [
        "Deploy ML-based anomaly detection (July 2025)",
        "Scale Kafka cluster to 48 partitions (August 2025)",
        "Implement cross-region failover (September 2025)",
        "Complete SOC 2 compliance audit (October 2025)",
    ]
    for item in next_items:
        p = tf6.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress for GUI agent
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
