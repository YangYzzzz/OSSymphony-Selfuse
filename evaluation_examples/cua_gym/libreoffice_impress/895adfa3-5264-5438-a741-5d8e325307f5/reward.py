"""
Reward Script: Insert a right-arrow connector shape between two boxes on slide 4
Task ID: impress_tm_060
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Right-arrow shape exists on slide 4
  Component 2 (0.35): Arrow positioned horizontally between Input and Output boxes
  Component 3 (0.30): Arrow vertically overlaps with the boxes and both boxes preserved
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_060'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def find_arrow_shapes(slide):
    """Find all arrow-type shapes on a slide by checking preset geometry."""
    arrows = []
    arrow_presets = {
        'rightArrow', 'leftArrow', 'upArrow', 'downArrow',
        'leftRightArrow', 'upDownArrow', 'bentArrow',
        'stripedRightArrow', 'notchedRightArrow', 'curvedRightArrow',
        'rightArrowCallout', 'blockArc',
    }
    for shape in slide.shapes:
        el = shape.element
        # Check for preset geometry
        prst_geom = el.find('.//' + qn('a:prstGeom'))
        if prst_geom is not None:
            prst = prst_geom.get('prst', '')
            if prst in arrow_presets or 'arrow' in prst.lower() or 'Arrow' in (shape.name or ''):
                arrows.append((shape, prst))
    return arrows


def find_labeled_box(slide, label):
    """Find a shape containing the given label text on a slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip().lower() == label.lower():
            # Must be a box-like shape (AUTO_SHAPE), not a text box placeholder
            return shape
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed, so slide 4 is index 3

    # Find Input and Output boxes (precondition gate)
    input_box = find_labeled_box(slide4, 'Input')
    output_box = find_labeled_box(slide4, 'Output')

    if input_box is None or output_box is None:
        print(f"FAIL: Cannot find Input and/or Output boxes on slide 4")
        print(f"  Input box found: {input_box is not None}")
        print(f"  Output box found: {output_box is not None}")
        print("REWARD: 0.0")
        return 0.0

    input_right = input_box.left + input_box.width
    output_left = output_box.left

    print(f"INFO: Input box right edge: {input_right}")
    print(f"INFO: Output box left edge: {output_left}")
    print(f"INFO: Input box vertical: top={input_box.top}, bottom={input_box.top + input_box.height}")
    print(f"INFO: Output box vertical: top={output_box.top}, bottom={output_box.top + output_box.height}")

    # Component 1: A right-arrow shape exists on slide 4 (0.35 points)
    # This is the core task: inserting an arrow shape
    try:
        arrow_shapes = find_arrow_shapes(slide4)
        if len(arrow_shapes) > 0:
            arrow_shape, arrow_prst = arrow_shapes[0]
            total_score += 0.35
            print(f"PASS: Component 1 — Right-arrow shape found on slide 4: "
                  f"name='{arrow_shape.name}', preset='{arrow_prst}' (0.35 pts)")
        else:
            # Also check for connector shapes (cxnSp elements)
            connector_shapes = [s for s in slide4.shapes
                                if s.element.tag.endswith('}cxnSp') or 'cxnSp' in s.element.tag]
            if len(connector_shapes) > 0:
                arrow_shape = connector_shapes[0]
                arrow_prst = 'connector'
                if arrow_shape is not None:
                    print(f"PASS: Component 1 — Connector shape found on slide 4: "
                          f"name='{arrow_shape.name}' (0.35 pts)")
                    total_score += 0.35
            else:
                print(f"FAIL: Component 1 — No arrow or connector shape found on slide 4")
                print(f"  Shapes on slide 4: {[s.name for s in slide4.shapes]}")
                arrow_shape = None
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        arrow_shape = None

    # Component 2: Arrow is positioned horizontally between Input and Output boxes (0.35 points)
    # The arrow's left edge should be >= Input right edge (or close)
    # The arrow's right edge should be <= Output left edge (or close)
    try:
        if arrow_shape is not None:
            arrow_left = arrow_shape.left
            arrow_right = arrow_shape.left + arrow_shape.width
            tolerance = 914400  # 1 inch tolerance

            # Arrow must be generally between the two boxes
            left_ok = arrow_left >= (input_right - tolerance)
            right_ok = arrow_right <= (output_left + tolerance)
            # Arrow must not be entirely to the left of input or right of output
            not_outside = arrow_right > input_right and arrow_left < output_left + output_box.width

            print(f"INFO: Arrow horizontal: left={arrow_left}, right={arrow_right}")
            print(f"INFO: left_ok={left_ok} (arrow_left {arrow_left} >= input_right-tol {input_right - tolerance})")
            print(f"INFO: right_ok={right_ok} (arrow_right {arrow_right} <= output_left+tol {output_left + tolerance})")

            if left_ok and right_ok and not_outside:
                print(f"PASS: Component 2 — Arrow positioned between Input and Output boxes (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 2 — Arrow not properly positioned between boxes")
        else:
            print(f"FAIL: Component 2 — No arrow shape to check position")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Arrow vertically overlaps with boxes AND both boxes preserved (0.30 points)
    # Arrow should be in the same vertical range as the Input/Output boxes
    try:
        if arrow_shape is not None:
            arrow_top = arrow_shape.top
            arrow_bottom = arrow_shape.top + arrow_shape.height

            # Input box vertical range
            input_top = input_box.top
            input_bottom = input_box.top + input_box.height

            # Check vertical overlap: arrow must overlap with the box vertical range
            overlap_top = max(arrow_top, input_top)
            overlap_bottom = min(arrow_bottom, input_bottom)
            has_vertical_overlap = overlap_bottom > overlap_top

            print(f"INFO: Arrow vertical: top={arrow_top}, bottom={arrow_bottom}")
            print(f"INFO: Vertical overlap: {has_vertical_overlap}")

            # Also verify both original boxes still exist with correct labels
            boxes_preserved = (input_box is not None and output_box is not None and
                               input_box.text.strip() == 'Input' and
                               output_box.text.strip() == 'Output')

            if has_vertical_overlap and boxes_preserved:
                print(f"PASS: Component 3 — Arrow vertically aligned with boxes, "
                      f"both Input/Output boxes preserved (0.30 pts)")
                total_score += 0.30
            elif has_vertical_overlap:
                print(f"PARTIAL: Component 3 — Vertical alignment OK but boxes modified (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 — Arrow not vertically aligned with boxes")
        else:
            print(f"FAIL: Component 3 — No arrow shape to check alignment")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
