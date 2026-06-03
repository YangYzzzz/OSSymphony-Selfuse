"""
Initial Setup: Process Flow presentation with arrow shape on slide 3
Task ID: impress_ma_062
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_062'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Process Flow Overview"
    slide1.placeholders[1].text = "Quarterly Operations Review - Q1 2025"

    # ── Slide 2: Process Introduction ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Process Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Content textbox
    content2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
    tf2 = content2.text_frame
    tf2.word_wrap = True
    lines = [
        "Our quarterly operations workflow consists of five key stages:",
        "1. Requirements Gathering - Collect stakeholder input and define scope",
        "2. Design & Planning - Create architecture diagrams and timeline",
        "3. Implementation - Execute development sprints with weekly reviews",
        "4. Quality Assurance - Automated testing and manual review cycles",
        "5. Deployment & Monitoring - Production rollout and performance tracking",
    ]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        r = p.runs[0]
        r.font.name = "Arial"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 3: Process Flow with Arrow ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Implementation Flow"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(32)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Process step boxes
    steps = ["Initiate", "Develop", "Test", "Review", "Deploy"]
    colors = [
        RGBColor(0x2E, 0x86, 0xC1),
        RGBColor(0x28, 0xB4, 0x63),
        RGBColor(0xF3, 0x9C, 0x12),
        RGBColor(0xE7, 0x4C, 0x3C),
        RGBColor(0x8E, 0x44, 0xAD),
    ]
    for i, (step, color) in enumerate(zip(steps, colors)):
        left = Inches(1.0 + i * 2.3)
        top = Inches(2.0)
        box = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.8), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        tf_box = box.text_frame
        tf_box.word_wrap = True
        p_box = tf_box.paragraphs[0]
        p_box.text = step
        p_box.alignment = PP_ALIGN.CENTER
        r_box = p_box.runs[0]
        r_box.font.name = "Arial"
        r_box.font.size = Pt(16)
        r_box.font.bold = True
        r_box.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Arrow shape on the LEFT side (the key shape for the task)
    arrow = slide3.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(0.5),   # x = 0.5in (left side)
        Inches(3.5),   # y = 3.5in
        Inches(1.5),   # width
        Inches(0.8),   # height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0xE6, 0x55, 0x00)
    arrow.line.color.rgb = RGBColor(0xCC, 0x44, 0x00)

    # Description text below arrow
    desc3 = slide3.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11), Inches(1.5))
    tf_desc = desc3.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "The arrow above represents the flow direction from project initiation through to final deployment."
    r_desc = p_desc.runs[0]
    r_desc.font.name = "Arial"
    r_desc.font.size = Pt(14)
    r_desc.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ── Slide 4: Timeline ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Project Timeline"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(32)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Timeline content
    timeline_items = [
        ("Week 1-2", "Requirements gathering and stakeholder interviews"),
        ("Week 3-4", "System design and architecture review"),
        ("Week 5-8", "Core development with bi-weekly sprint demos"),
        ("Week 9-10", "Integration testing and UAT cycles"),
        ("Week 11-12", "Staged rollout and performance monitoring"),
    ]
    for i, (period, desc) in enumerate(timeline_items):
        y = Inches(1.8 + i * 1.0)
        # Period label
        lbl = slide4.shapes.add_textbox(Inches(1.0), y, Inches(2.5), Inches(0.6))
        tf_lbl = lbl.text_frame
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = period
        r_lbl = p_lbl.runs[0]
        r_lbl.font.name = "Arial"
        r_lbl.font.size = Pt(16)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = RGBColor(0x2E, 0x86, 0xC1)

        # Description
        dtx = slide4.shapes.add_textbox(Inches(3.8), y, Inches(8), Inches(0.6))
        tf_dtx = dtx.text_frame
        p_dtx = tf_dtx.paragraphs[0]
        p_dtx.text = desc
        r_dtx = p_dtx.runs[0]
        r_dtx.font.name = "Arial"
        r_dtx.font.size = Pt(14)
        r_dtx.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 5: Summary ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    txBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Takeaways"
    p5.alignment = PP_ALIGN.LEFT
    run5 = p5.runs[0]
    run5.font.name = "Arial"
    run5.font.size = Pt(32)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    takeaways = [
        "Cross-functional collaboration reduces cycle time by 30%",
        "Automated testing catches 85% of regressions before production",
        "Weekly stakeholder touchpoints maintain alignment throughout the project",
        "Post-deployment monitoring dashboards enable proactive issue resolution",
    ]
    content5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(4.5))
    tf_c5 = content5.text_frame
    tf_c5.word_wrap = True
    for i, item in enumerate(takeaways):
        if i == 0:
            p_c = tf_c5.paragraphs[0]
        else:
            p_c = tf_c5.add_paragraph()
        p_c.text = item
        r_c = p_c.runs[0]
        r_c.font.name = "Arial"
        r_c.font.size = Pt(18)
        r_c.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
