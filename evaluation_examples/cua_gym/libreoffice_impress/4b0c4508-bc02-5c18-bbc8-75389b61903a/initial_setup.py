"""
Initial Setup: Create a chemistry lab presentation with 6 slides.
Slide 4 titled 'Results' has NO notes.
Task ID: impress_ndo_027
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_027'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_name="Calibri", font_size=18,
                      bold=False, color=None, alignment=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Chemistry Lab Report"
    slide1.placeholders[1].text = "Thermodynamic Reaction Analysis\nDr. Elena Martinez | Spring 2025"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Objective: Investigate the exothermic reaction between sodium hydroxide and hydrochloric acid"
    p2a = body2.add_paragraph()
    p2a.text = "Hypothesis: The reaction will produce a measurable temperature increase proportional to concentration"
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "Background: Neutralization reactions release energy as heat, governed by enthalpy changes"
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "Equipment: Calorimeter, thermometer, graduated cylinders, safety goggles"
    p2c.level = 0

    # --- Slide 3: Methodology ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Methodology"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "1. Measured 50 mL of 1.0 M NaOH into the calorimeter"
    p3a = body3.add_paragraph()
    p3a.text = "2. Recorded initial temperature of both solutions (22\u00b0C)"
    p3a.level = 0
    p3b = body3.add_paragraph()
    p3b.text = "3. Added 50 mL of 1.0 M HCl to the calorimeter"
    p3b.level = 0
    p3c = body3.add_paragraph()
    p3c.text = "4. Monitored temperature every 5 seconds for 3 minutes"
    p3c.level = 0
    p3d = body3.add_paragraph()
    p3d.text = "5. Repeated the experiment three times for statistical validity"
    p3d.level = 0

    # --- Slide 4: Results (NO NOTES - this is the target slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Results"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Peak temperature reached: 34\u00b0C (from 22\u00b0C baseline)"
    p4a = body4.add_paragraph()
    p4a.text = "Average temperature increase across trials: 12.1\u00b0C (\u00b10.3\u00b0C)"
    p4a.level = 0
    p4b = body4.add_paragraph()
    p4b.text = "Atmospheric pressure: 1.013 atm (constant throughout)"
    p4b.level = 0
    p4c = body4.add_paragraph()
    p4c.text = "Mean reaction time to peak: 45.2 seconds"
    p4c.level = 0
    p4d = body4.add_paragraph()
    p4d.text = "Enthalpy change: -57.1 kJ/mol"
    p4d.level = 0
    # NOTE: Do NOT access slide4.notes_slide -- that would create notes

    # --- Slide 5: Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Analysis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "The measured enthalpy aligns with the theoretical value of -57.3 kJ/mol"
    p5a = body5.add_paragraph()
    p5a.text = "Temperature increase was consistent across all three trials (low standard deviation)"
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Constant pressure confirms the reaction followed isobaric conditions"
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Minor heat loss to surroundings accounts for 0.2 kJ/mol deviation"
    p5c.level = 0

    # --- Slide 6: Conclusion ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusion"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "The exothermic neutralization reaction produced predictable, reproducible results"
    p6a = body6.add_paragraph()
    p6a.text = "Experimental enthalpy closely matched theoretical predictions"
    p6a.level = 0
    p6b = body6.add_paragraph()
    p6b.text = "Future work: Investigate concentration effects on reaction rate and heat output"
    p6b.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
