"""
Reward Script: Add structured notes to slide 4 of ChemistryLab.pptx
Task ID: impress_ndo_027
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Bold header 'Experiment Results' in slide 4 notes
  Component 2 (0.25): Bullet '- Temperature increased by 12 deg C' (non-bold)
  Component 3 (0.25): Bullet '- Pressure remained constant at 1 atm' (non-bold)
  Component 4 (0.20): Bullet '- Reaction time: 45 seconds' (non-bold)
"""

import os

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_027'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_notes_paragraphs(slide):
    """Get notes paragraphs from a slide, or None if no notes exist."""
    # Check if notes relationship exists without creating one
    for rel in slide.part.rels.values():
        if rel.reltype == RT.NOTES_SLIDE:
            return slide.notes_slide.notes_text_frame.paragraphs
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Precondition: notes must exist on slide 4
    paras = get_notes_paragraphs(slide)
    if paras is None:
        print("FAIL: Slide 4 has no notes slide attached")
        print("REWARD: 0.0")
        return 0.0

    # Collect non-empty paragraphs from notes
    notes_paras = []
    for p in paras:
        text = p.text.strip()
        if text:
            notes_paras.append(p)

    if len(notes_paras) == 0:
        print("FAIL: Slide 4 notes are empty (no non-empty paragraphs)")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Bold header 'Experiment Results' (0.30 points)
    try:
        first_para = notes_paras[0]
        header_text = first_para.text.strip()
        # Check text content
        header_text_ok = (header_text == "Experiment Results")
        # Check bold formatting on all runs
        runs = [r for r in first_para.runs if (r.text or "").strip()]
        header_bold_ok = False
        if runs:
            header_bold_ok = all(r.font.bold is True for r in runs)

        if header_text_ok and header_bold_ok:
            print(f"PASS: Component 1 -- Bold header 'Experiment Results' found (0.30 pts)")
            total_score += 0.30
        elif header_text_ok and not header_bold_ok:
            # Partial: text correct but not bold -- give small partial credit
            print(f"FAIL: Component 1 -- Header text correct but not bold (runs bold: {[r.font.bold for r in runs]})")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 -- Expected 'Experiment Results', found: {header_text!r}, bold: {header_bold_ok}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Expected bullet lines
    expected_bullets = [
        "- Temperature increased by 12\u00b0C",
        "- Pressure remained constant at 1 atm",
        "- Reaction time: 45 seconds",
    ]
    bullet_weights = [0.25, 0.25, 0.20]

    # Component 2-4: Bullet point lines (non-bold)
    for idx, (expected_text, weight) in enumerate(zip(expected_bullets, bullet_weights)):
        comp_num = idx + 2
        try:
            # Look for this bullet in notes paragraphs (skip the header)
            found = False
            for para in notes_paras[1:]:
                para_text = para.text.strip()
                if para_text == expected_text:
                    # Check that runs are NOT bold
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    if runs:
                        all_not_bold = all(
                            (r.font.bold is None or r.font.bold is False) for r in runs
                        )
                    else:
                        # Text present but no runs (direct text) -- accept
                        all_not_bold = True

                    if all_not_bold:
                        print(f"PASS: Component {comp_num} -- Bullet '{expected_text}' found, non-bold ({weight} pts)")
                        total_score += weight
                    else:
                        print(f"FAIL: Component {comp_num} -- Bullet '{expected_text}' found but is bold")
                        total_score += weight * 0.4  # partial: text ok but formatting wrong
                    found = True
                    break

            if not found:
                # Try fuzzy match (degree symbol variants, whitespace)
                for para in notes_paras[1:]:
                    para_text = para.text.strip()
                    # Normalize: replace various degree symbols and whitespace
                    norm_expected = expected_text.replace('\u00b0', 'deg').replace('  ', ' ')
                    norm_actual = para_text.replace('\u00b0', 'deg').replace('  ', ' ')
                    if norm_expected.lower() == norm_actual.lower():
                        print(f"PASS: Component {comp_num} -- Bullet '{expected_text}' found (fuzzy match) ({weight} pts)")
                        total_score += weight
                        found = True
                        break

                if not found:
                    actual_bullets = [p.text.strip() for p in notes_paras[1:]]
                    print(f"FAIL: Component {comp_num} -- Bullet '{expected_text}' not found. Actual bullets: {actual_bullets}")

        except Exception as e:
            print(f"ERROR: Component {comp_num} -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
