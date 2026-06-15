"""
Reward Script: Verify statistical analysis slides 5-8 in psychology experiment presentation
Task ID: impress_stu_087
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide 5 has a descriptive statistics table with Mean, SD, N
  Component 2 (0.15): Slide 6 has a chart (box plot / bar chart for group comparison)
  Component 3 (0.20): Slide 7 has a chart + text box with t-test result string
  Component 4 (0.20): Slide 8 has 4 bullet points + highlighted key takeaway box
  Component 5 (0.15): Slide 8 highlighted box has yellow background and dark border
  Component 6 (0.10): Consistent #2C3E50 title color across slides 5-8
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_087'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 5 has a descriptive statistics table (0.20 points)
    # Initial: slide 5 has only a title placeholder; golden adds a 4x4 table
    try:
        slide5 = prs.slides[4]  # 0-indexed
        table_found = False
        table_has_stats = False
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_found = True
                # Check for expected headers: Mean, SD, N
                header_texts = []
                for c in range(len(table.columns)):
                    header_texts.append(table.cell(0, c).text.strip().lower())
                has_mean = any('mean' in h for h in header_texts)
                has_sd = any('sd' in h for h in header_texts)
                has_n = any(h == 'n' for h in header_texts)
                # Check for control/experimental group rows
                all_cell_text = ''
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        all_cell_text += table.cell(r, c).text.lower() + ' '
                has_control = 'control' in all_cell_text
                has_experimental = 'experimental' in all_cell_text
                if has_mean and has_sd and has_n and has_control and has_experimental:
                    table_has_stats = True
                    print(f"PASS: Component 1 — Slide 5 table has Mean/SD/N headers and control/experimental rows (0.20 pts)")
                    total_score += 0.20
                elif table_found:
                    print(f"PARTIAL: Component 1 — Table found but missing expected headers/rows. Headers: {header_texts}")
                    total_score += 0.10
                break
        if not table_found:
            print(f"FAIL: Component 1 — No table found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 6 has a chart (group comparison / box plot) (0.15 points)
    # Initial: slide 6 has only a title; golden adds a chart
    try:
        slide6 = prs.slides[5]
        chart_found = False
        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_found = True
                print(f"PASS: Component 2 — Slide 6 has a chart (type={shape.chart.chart_type}) (0.15 pts)")
                total_score += 0.15
                break
        if not chart_found:
            print(f"FAIL: Component 2 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 7 has a chart + text with t-test results (0.20 points)
    # Initial: slide 7 has only a title; golden adds a chart and a text box with stats
    try:
        slide7 = prs.slides[6]
        chart_found = False
        ttest_text_found = False
        for shape in slide7.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_found = True
            if shape.has_text_frame:
                full_text = ''
                for para in shape.text_frame.paragraphs:
                    full_text += para.text + ' '
                # Check for the t-test result string components
                ft_lower = full_text.lower()
                if 't(48)' in full_text and 'p <' in full_text and 'd =' in full_text:
                    ttest_text_found = True
                elif 't(48)' in full_text or ('3.21' in full_text and '0.92' in full_text):
                    ttest_text_found = True

        comp3_score = 0.0
        if chart_found and ttest_text_found:
            comp3_score = 0.20
            print(f"PASS: Component 3 — Slide 7 has chart + t-test result text (0.20 pts)")
        elif chart_found:
            comp3_score = 0.10
            print(f"PARTIAL: Component 3 — Slide 7 has chart but missing t-test text (0.10 pts)")
        elif ttest_text_found:
            comp3_score = 0.10
            print(f"PARTIAL: Component 3 — Slide 7 has t-test text but missing chart (0.10 pts)")
        else:
            print(f"FAIL: Component 3 — Slide 7 missing both chart and t-test text")
        total_score += comp3_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 8 has bullet points (>= 3) + separate key takeaway text box (0.20 points)
    # Initial: slide 8 has only a title; golden adds 2 text boxes
    try:
        slide8 = prs.slides[7]
        # Count non-title text boxes
        text_boxes = []
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
                text_boxes.append(shape)

        bullet_box = None
        takeaway_box = None
        for tb in text_boxes:
            full_text = ''
            para_count = 0
            for para in tb.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    para_count += 1
                    full_text += txt + ' '
            ft_lower = full_text.lower()
            if 'takeaway' in ft_lower or 'key' in ft_lower:
                takeaway_box = tb
            elif para_count >= 3:
                bullet_box = tb

        comp4_score = 0.0
        if bullet_box is not None and takeaway_box is not None:
            # Count bullet paragraphs in the bullet box
            bullet_count = sum(1 for p in bullet_box.text_frame.paragraphs if p.text.strip())
            if bullet_count >= 4:
                comp4_score = 0.20
                print(f"PASS: Component 4 — Slide 8 has {bullet_count} bullet points + key takeaway box (0.20 pts)")
            else:
                comp4_score = 0.15
                print(f"PARTIAL: Component 4 — Slide 8 has {bullet_count} bullet points (need 4) + key takeaway box (0.15 pts)")
        elif bullet_box is not None or takeaway_box is not None:
            comp4_score = 0.10
            has_what = 'bullet box' if bullet_box else 'takeaway box'
            print(f"PARTIAL: Component 4 — Slide 8 has {has_what} but missing the other (0.10 pts)")
        else:
            # Check if there's at least some content beyond the title
            if len(text_boxes) > 0:
                comp4_score = 0.05
                print(f"PARTIAL: Component 4 — Slide 8 has text content but cannot identify bullet/takeaway structure (0.05 pts)")
            else:
                print(f"FAIL: Component 4 — Slide 8 has no text boxes beyond the title")
        total_score += comp4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 8 highlighted box has yellow background and dark border (0.15 points)
    # Initial: no such box exists; golden has a box with fill FFF176 and border 2C3E50
    try:
        slide8 = prs.slides[7]
        yellow_fill_found = False
        dark_border_found = False
        for shape in slide8.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.has_text_frame:
                # Check for yellow-ish fill
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    try:
                        rgb = fill.fore_color.rgb
                        r_val = int(str(rgb)[0:2], 16)
                        g_val = int(str(rgb)[2:4], 16)
                        b_val = int(str(rgb)[4:6], 16)
                        # Yellow: high R, high G, low-ish B
                        if r_val > 200 and g_val > 200 and b_val < 150:
                            yellow_fill_found = True
                            print(f"  Found yellow fill: {rgb}")
                    except Exception:
                        pass
                # Check for dark border
                line = shape.line
                if line.fill.type is not None and line.fill.type == 1:
                    try:
                        line_rgb = line.color.rgb
                        lr = int(str(line_rgb)[0:2], 16)
                        lg = int(str(line_rgb)[2:4], 16)
                        lb = int(str(line_rgb)[4:6], 16)
                        # Dark: all channels < 100, or matches #2C3E50
                        if (lr < 100 and lg < 100 and lb < 100) or str(line_rgb) == '2C3E50':
                            dark_border_found = True
                            print(f"  Found dark border: {line_rgb}")
                    except Exception:
                        pass

        comp5_score = 0.0
        if yellow_fill_found and dark_border_found:
            comp5_score = 0.15
            print(f"PASS: Component 5 — Highlighted box has yellow fill + dark border (0.15 pts)")
        elif yellow_fill_found:
            comp5_score = 0.10
            print(f"PARTIAL: Component 5 — Yellow fill found but no dark border (0.10 pts)")
        elif dark_border_found:
            comp5_score = 0.05
            print(f"PARTIAL: Component 5 — Dark border found but no yellow fill (0.05 pts)")
        else:
            print(f"FAIL: Component 5 — No highlighted box with yellow fill and dark border found on slide 8")
        total_score += comp5_score
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Consistent #2C3E50 title color across slides 5-8 (0.10 points)
    # Note: This is true in initial too for the title text. However, in initial slides 5-8 ONLY
    # have a title and nothing else. We check that titles still maintain #2C3E50 AFTER content
    # was added (i.e., the task didn't break the formatting). This component is a compound check:
    # title color is #2C3E50 AND the slide has additional content (not title-only).
    try:
        correct_titles = 0
        for slide_idx in [4, 5, 6, 7]:  # slides 5-8
            slide = prs.slides[slide_idx]
            # Check this slide has content beyond the title
            non_title_shapes = [s for s in slide.shapes
                                if s.shape_type != 14]  # 14 = PLACEHOLDER
            if len(non_title_shapes) == 0:
                # No content added — this is the initial state, don't award points
                continue
            # Check title color
            for shape in slide.shapes:
                if shape.name and 'title' in shape.name.lower() and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None and str(run.font.color.rgb) == '2C3E50':
                                    correct_titles += 1
                            except:
                                pass
                    break  # Only check title shape

        if correct_titles >= 4:
            print(f"PASS: Component 6 — All 4 slides (5-8) have #2C3E50 title color with content (0.10 pts)")
            total_score += 0.10
        elif correct_titles >= 2:
            partial = round(0.10 * correct_titles / 4, 2)
            print(f"PARTIAL: Component 6 — {correct_titles}/4 slides have correct title color with content ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 6 — Only {correct_titles}/4 slides have #2C3E50 title color with content")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
