"""
Initial Setup: Psychology experiment presentation with slides 5-8 having titles only.
Task ID: impress_stu_087
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_087'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

TITLE_COLOR = RGBColor(0x2C, 0x3E, 0x50)


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title(slide, text):
    title_shape = slide.shapes.title
    title_shape.text = text
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR


def add_body_text(slide, text, left=Inches(0.8), top=Inches(1.8),
                  width=Inches(8.4), height=Inches(4.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Layout 0 = Title Slide, Layout 5 = Title Only
    layout_title = prs.slide_layouts[0]
    layout_title_only = prs.slide_layouts[5]

    # --- Slide 1: Title Slide ---
    s1 = prs.slides.add_slide(layout_title)
    s1.shapes.title.text = "Cognitive Load and Decision-Making Quality"
    for run in s1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = TITLE_COLOR
    s1.placeholders[1].text = (
        "A Randomized Controlled Experiment\n"
        "Dr. Elena Vasquez & Prof. James Whitfield\n"
        "Department of Psychology, Lakewood University"
    )

    # --- Slide 2: Introduction ---
    s2 = prs.slides.add_slide(layout_title_only)
    set_title(s2, "Introduction")
    add_body_text(s2,
        "Cognitive load theory (Sweller, 1988) suggests that working memory limitations "
        "directly impact decision-making quality. Previous research has shown mixed results "
        "regarding the threshold at which cognitive load begins to degrade performance.\n\n"
        "This study examines whether participants under high cognitive load (dual-task paradigm) "
        "make significantly worse decisions compared to a control group performing the same "
        "decision tasks without additional cognitive demands."
    )

    # --- Slide 3: Research Questions ---
    s3 = prs.slides.add_slide(layout_title_only)
    set_title(s3, "Research Questions & Hypotheses")
    add_body_text(s3,
        "RQ1: Does increased cognitive load reduce decision accuracy?\n\n"
        "RQ2: Is there a measurable difference in response time between groups?\n\n"
        "H1: Participants in the experimental (high-load) group will show significantly "
        "lower decision accuracy than the control group.\n\n"
        "H2: Response times will be significantly longer in the high-load condition."
    )

    # --- Slide 4: Methodology ---
    s4 = prs.slides.add_slide(layout_title_only)
    set_title(s4, "Methodology")
    add_body_text(s4,
        "Participants: 50 undergraduate students (25 per group), ages 18-24\n"
        "Design: Between-subjects, randomized assignment\n\n"
        "Control Group: Standard decision-making task battery (Iowa Gambling Task)\n"
        "Experimental Group: Same task battery + concurrent auditory n-back task\n\n"
        "Measures:\n"
        "  - Decision accuracy (% correct optimal choices)\n"
        "  - Response time (ms per decision)\n"
        "  - Subjective cognitive load (NASA-TLX scale)"
    )

    # --- Slide 5: Descriptive Statistics (TITLE ONLY - task target) ---
    s5 = prs.slides.add_slide(layout_title_only)
    set_title(s5, "Descriptive Statistics")

    # --- Slide 6: Group Comparison (TITLE ONLY - task target) ---
    s6 = prs.slides.add_slide(layout_title_only)
    set_title(s6, "Group Comparison")

    # --- Slide 7: Statistical Significance (TITLE ONLY - task target) ---
    s7 = prs.slides.add_slide(layout_title_only)
    set_title(s7, "Statistical Significance")

    # --- Slide 8: Results Summary (TITLE ONLY - task target) ---
    s8 = prs.slides.add_slide(layout_title_only)
    set_title(s8, "Results Summary")

    # --- Slide 9: Discussion ---
    s9 = prs.slides.add_slide(layout_title_only)
    set_title(s9, "Discussion")
    add_body_text(s9,
        "The findings support cognitive load theory predictions. Participants under "
        "dual-task conditions demonstrated significantly reduced decision accuracy, "
        "consistent with Sweller's framework.\n\n"
        "The large effect size (d = 0.92) suggests practical significance beyond "
        "statistical significance. These results have implications for real-world "
        "settings where multitasking during critical decisions is common."
    )

    # --- Slide 10: References ---
    s10 = prs.slides.add_slide(layout_title_only)
    set_title(s10, "References")
    add_body_text(s10,
        "Bechara, A., et al. (1994). Insensitivity to future consequences following "
        "damage to human prefrontal cortex. Cognition, 50(1-3), 7-15.\n\n"
        "Hart, S. G., & Staveland, L. E. (1988). Development of NASA-TLX. "
        "Advances in Psychology, 52, 139-183.\n\n"
        "Sweller, J. (1988). Cognitive load during problem solving. "
        "Cognitive Science, 12(2), 257-285.\n\n"
        "Kahneman, D. (2011). Thinking, Fast and Slow. Farrar, Straus and Giroux.",
        top=Inches(1.6)
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
