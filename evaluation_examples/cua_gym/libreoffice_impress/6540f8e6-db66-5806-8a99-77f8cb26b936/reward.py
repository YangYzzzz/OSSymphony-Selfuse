"""
Reward Script: Move images to top of slides 2 and 5, underline body text on slides 3, 4, 6
Task ID: osworld_impress_image_top_underline_text_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Picture on slide 2 moved to top (top < 15% of slide height)
  Component 2 (0.30): Picture on slide 5 moved to top (top < 15% of slide height)
  Component 3 (0.13): All body text runs on slide 3 are underlined
  Component 4 (0.13): All body text runs on slide 4 are underlined
  Component 5 (0.14): All body text runs on slide 6 are underlined
  Total: 1.00
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_007'

# "Top of slide" threshold: picture top must be < 15% of slide height.
# Initial golden top = 457200 EMU (6.7%); initial non-moved top = 1600200 EMU (23.3%).
TOP_THRESHOLD_PCT = 0.15

# Body textbox name to look for on slides 3, 4, 6
BODY_TEXTBOX_NAME = 'TextBox 3'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height
    top_threshold = int(slide_height * TOP_THRESHOLD_PCT)
    total_score = 0.0

    print(f"Slide count: {len(prs.slides)}, slide_height: {slide_height}, top_threshold: {top_threshold}")

    # ----------------------------------------------------------------
    # Component 1: Picture on slide 2 moved to top (0.30 points)
    # ----------------------------------------------------------------
    # Task: move infographic image on slide 2 to the top of slide 2.
    # "Top" means picture.top < TOP_THRESHOLD_PCT * slide_height.
    # On initial: top=1600200 (23.3%), on golden: top=457200 (6.7%).
    try:
        slide2 = prs.slides[1]
        picture_on_slide2 = None
        for shape in slide2.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_on_slide2 = shape
                break

        if picture_on_slide2 is None:
            print("FAIL: Component 1 — no picture shape found on slide 2")
        else:
            actual_top = picture_on_slide2.top
            if actual_top < top_threshold:
                print(f"PASS: Component 1 — picture on slide 2 at top={actual_top} < threshold {top_threshold} (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — picture on slide 2 at top={actual_top}, expected < {top_threshold}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ----------------------------------------------------------------
    # Component 2: Picture on slide 5 moved to top (0.30 points)
    # ----------------------------------------------------------------
    # Task: move chart/image on slide 5 to the top of slide 5.
    try:
        slide5 = prs.slides[4]
        picture_on_slide5 = None
        for shape in slide5.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_on_slide5 = shape
                break

        if picture_on_slide5 is None:
            print("FAIL: Component 2 — no picture shape found on slide 5")
        else:
            actual_top = picture_on_slide5.top
            if actual_top < top_threshold:
                print(f"PASS: Component 2 — picture on slide 5 at top={actual_top} < threshold {top_threshold} (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 2 — picture on slide 5 at top={actual_top}, expected < {top_threshold}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ----------------------------------------------------------------
    # Component 3: All body text runs on slide 3 are underlined (0.13 points)
    # ----------------------------------------------------------------
    # Task: underline all body text on slide 3.
    # Body textbox = 'TextBox 3' (the large multi-paragraph content box).
    try:
        slide3 = prs.slides[2]
        body_shape_3 = None
        for shape in slide3.shapes:
            if shape.name == BODY_TEXTBOX_NAME and shape.has_text_frame:
                body_shape_3 = shape
                break

        if body_shape_3 is None:
            print(f"FAIL: Component 3 — '{BODY_TEXTBOX_NAME}' not found on slide 3")
        else:
            runs_found = []
            for para in body_shape_3.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs_found.append(run)

            if not runs_found:
                print("FAIL: Component 3 — no non-empty runs found in body textbox on slide 3")
            else:
                all_underlined = all(run.font.underline is True for run in runs_found)
                if all_underlined:
                    print(f"PASS: Component 3 — all {len(runs_found)} body runs on slide 3 are underlined (0.13 pts)")
                    total_score += 0.13
                else:
                    not_underlined = [run.text[:30] for run in runs_found if run.font.underline is not True]
                    print(f"FAIL: Component 3 — {len(not_underlined)} body runs on slide 3 not underlined: {not_underlined[:3]}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ----------------------------------------------------------------
    # Component 4: All body text runs on slide 4 are underlined (0.13 points)
    # ----------------------------------------------------------------
    # Task: underline all body text on slide 4.
    try:
        slide4 = prs.slides[3]
        body_shape_4 = None
        for shape in slide4.shapes:
            if shape.name == BODY_TEXTBOX_NAME and shape.has_text_frame:
                body_shape_4 = shape
                break

        if body_shape_4 is None:
            print(f"FAIL: Component 4 — '{BODY_TEXTBOX_NAME}' not found on slide 4")
        else:
            runs_found = []
            for para in body_shape_4.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs_found.append(run)

            if not runs_found:
                print("FAIL: Component 4 — no non-empty runs found in body textbox on slide 4")
            else:
                all_underlined = all(run.font.underline is True for run in runs_found)
                if all_underlined:
                    print(f"PASS: Component 4 — all {len(runs_found)} body runs on slide 4 are underlined (0.13 pts)")
                    total_score += 0.13
                else:
                    not_underlined = [run.text[:30] for run in runs_found if run.font.underline is not True]
                    print(f"FAIL: Component 4 — {len(not_underlined)} body runs on slide 4 not underlined: {not_underlined[:3]}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ----------------------------------------------------------------
    # Component 5: All body text runs on slide 6 are underlined (0.14 points)
    # ----------------------------------------------------------------
    # Task: underline all body text on slide 6.
    try:
        slide6 = prs.slides[5]
        body_shape_6 = None
        for shape in slide6.shapes:
            if shape.name == BODY_TEXTBOX_NAME and shape.has_text_frame:
                body_shape_6 = shape
                break

        if body_shape_6 is None:
            print(f"FAIL: Component 5 — '{BODY_TEXTBOX_NAME}' not found on slide 6")
        else:
            runs_found = []
            for para in body_shape_6.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        runs_found.append(run)

            if not runs_found:
                print("FAIL: Component 5 — no non-empty runs found in body textbox on slide 6")
            else:
                all_underlined = all(run.font.underline is True for run in runs_found)
                if all_underlined:
                    print(f"PASS: Component 5 — all {len(runs_found)} body runs on slide 6 are underlined (0.14 pts)")
                    total_score += 0.14
                else:
                    not_underlined = [run.text[:30] for run in runs_found if run.font.underline is not True]
                    print(f"FAIL: Component 5 — {len(not_underlined)} body runs on slide 6 not underlined: {not_underlined[:3]}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Run verification on the canonical artifact path in the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
