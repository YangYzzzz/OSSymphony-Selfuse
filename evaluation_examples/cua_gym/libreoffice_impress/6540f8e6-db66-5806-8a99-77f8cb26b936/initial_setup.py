"""
Initial Setup: 7-slide business intelligence deck
Task ID: osworld_impress_image_top_underline_text_007
Domain: libreoffice_impress

Creates a presentation with:
- Slide 2: infographic image placed in CENTER (not top)
- Slide 5: chart image placed in CENTER (not top)
- Slides 3, 4, 6: body text that is NOT underlined
"""

import os
import io
import base64
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_image_top_underline_text_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_infographic_image():
    """Create a business infographic image and return bytes."""
    img = Image.new('RGB', (600, 400), color='#1565C0')
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 590, 390], outline='white', width=3)
    # Title area
    draw.rectangle([10, 10, 590, 70], fill='#0D47A1')
    draw.text((300, 40), 'Business Intelligence Dashboard', fill='white', anchor='mm')
    # Bar chart
    bars = [(80, 160, '#4CAF50', 'Q1'), (200, 200, '#2196F3', 'Q2'),
            (320, 240, '#FF9800', 'Q3'), (440, 180, '#E91E63', 'Q4')]
    for x, h, c, label in bars:
        draw.rectangle([x, 300 - h, x + 70, 300], fill=c)
        draw.text((x + 35, 310), label, fill='white', anchor='mm')
    draw.text((300, 360), 'Annual Revenue Performance 2024', fill='#FFD700', anchor='mm')
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()


def make_chart_image():
    """Create a business chart image and return bytes."""
    img = Image.new('RGB', (600, 400), color='#263238')
    draw = ImageDraw.Draw(img)
    draw.rectangle([10, 10, 590, 390], outline='#78909C', width=2)
    # Title
    draw.rectangle([10, 10, 590, 60], fill='#1C313A')
    draw.text((300, 35), 'Market Share Analysis 2024', fill='white', anchor='mm')
    # Pie-like segments (simple rectangles for segments)
    segments = [
        (30, 80, 200, 350, '#42A5F5', 'Product A\n38%'),
        (220, 80, 370, 250, '#66BB6A', 'Product B\n25%'),
        (390, 80, 570, 200, '#FFA726', 'Product C\n20%'),
        (220, 270, 370, 350, '#EF5350', 'Product D\n17%'),
    ]
    for x1, y1, x2, y2, color, label in segments:
        draw.rectangle([x1, y1, x2, y2], fill=color, outline='white', width=2)
        cx, cy = (x1 + x2) // 2, (y1 + y2) // 2
        draw.text((cx, cy), label, fill='white', anchor='mm')
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = 'Business Intelligence Report 2024'
    slide1.placeholders[1].text = 'Prepared by the Analytics Team\nQ4 Performance Summary'

    # ---- Slide 2: Infographic Image in CENTER (not top) ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Add title text box
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = 'Annual Revenue Dashboard'
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.size = Pt(28)
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Place infographic image in CENTER of slide
    infographic_bytes = make_infographic_image()
    img_stream = io.BytesIO(infographic_bytes)
    img_w = Inches(6)
    img_h = Inches(4)
    img_left = (slide_w - img_w) // 2   # horizontally centered
    img_top = (slide_h - img_h) // 2    # vertically centered (NOT top)
    slide2.shapes.add_picture(img_stream, img_left, img_top, img_w, img_h)

    # ---- Slide 3: KPI Summary with body text (NOT underlined) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf3t = title3.text_frame
    tf3t.paragraphs[0].text = 'Key Performance Indicators'
    tf3t.paragraphs[0].runs[0].font.bold = True
    tf3t.paragraphs[0].runs[0].font.size = Pt(28)
    # Body textbox (NOT underlined)
    body3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5))
    tf3 = body3.text_frame
    tf3.word_wrap = True
    lines3 = [
        'Total Revenue grew 24% year-over-year, reaching $4.2M in Q4 2024.',
        'Customer acquisition cost decreased by 18% through optimized ad targeting.',
        'Net Promoter Score improved from 52 to 67, indicating strong customer satisfaction.',
        'Product return rate dropped to 3.1%, lowest in company history.',
        'Operating margin expanded to 31.5%, driven by cost efficiency initiatives.',
    ]
    p = tf3.paragraphs[0]
    p.text = lines3[0]
    run = p.runs[0]
    run.font.size = Pt(16)
    run.font.underline = False  # explicitly NOT underlined
    for line in lines3[1:]:
        para = tf3.add_paragraph()
        para.text = line
        r = para.runs[0]
        r.font.size = Pt(16)
        r.font.underline = False  # explicitly NOT underlined

    # ---- Slide 4: Regional Analysis with body text (NOT underlined) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf4t = title4.text_frame
    tf4t.paragraphs[0].text = 'Regional Sales Analysis'
    tf4t.paragraphs[0].runs[0].font.bold = True
    tf4t.paragraphs[0].runs[0].font.size = Pt(28)
    # Body textbox (NOT underlined)
    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5))
    tf4 = body4.text_frame
    tf4.word_wrap = True
    lines4 = [
        'North America: $1.8M revenue, up 31% from prior year with 14 new enterprise clients.',
        'Europe & Middle East: $1.2M revenue, steady growth despite macroeconomic headwinds.',
        'Asia Pacific: $0.9M revenue, fastest growing region at 42% YoY expansion.',
        'Latin America: $0.3M revenue, emerging market with high long-term potential.',
        'Cross-regional collaboration resulted in 6 joint deals worth $450K combined.',
    ]
    p4 = tf4.paragraphs[0]
    p4.text = lines4[0]
    r4 = p4.runs[0]
    r4.font.size = Pt(16)
    r4.font.underline = False  # explicitly NOT underlined
    for line in lines4[1:]:
        para = tf4.add_paragraph()
        para.text = line
        r = para.runs[0]
        r.font.size = Pt(16)
        r.font.underline = False  # explicitly NOT underlined

    # ---- Slide 5: Market Share Chart in CENTER (not top) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf5t = title5.text_frame
    tf5t.paragraphs[0].text = 'Market Share Analysis'
    tf5t.paragraphs[0].runs[0].font.bold = True
    tf5t.paragraphs[0].runs[0].font.size = Pt(28)
    tf5t.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

    # Place chart image in CENTER of slide
    chart_bytes = make_chart_image()
    chart_stream = io.BytesIO(chart_bytes)
    chart_w = Inches(6)
    chart_h = Inches(4)
    chart_left = (slide_w - chart_w) // 2    # horizontally centered
    chart_top = (slide_h - chart_h) // 2     # vertically centered (NOT top)
    slide5.shapes.add_picture(chart_stream, chart_left, chart_top, chart_w, chart_h)

    # ---- Slide 6: Strategic Initiatives with body text (NOT underlined) ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf6t = title6.text_frame
    tf6t.paragraphs[0].text = 'Strategic Initiatives for 2025'
    tf6t.paragraphs[0].runs[0].font.bold = True
    tf6t.paragraphs[0].runs[0].font.size = Pt(28)
    # Body textbox (NOT underlined)
    body6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.5), Inches(5))
    tf6 = body6.text_frame
    tf6.word_wrap = True
    lines6 = [
        'Launch AI-powered analytics platform to reduce manual reporting by 60%.',
        'Expand sales team by 25 FTEs across North America and Asia Pacific regions.',
        'Implement subscription-based pricing model targeting mid-market segment.',
        'Achieve ISO 27001 certification to unlock government and healthcare contracts.',
        'Establish strategic partnerships with 3 system integrators in EMEA territory.',
    ]
    p6 = tf6.paragraphs[0]
    p6.text = lines6[0]
    r6 = p6.runs[0]
    r6.font.size = Pt(16)
    r6.font.underline = False  # explicitly NOT underlined
    for line in lines6[1:]:
        para = tf6.add_paragraph()
        para.text = line
        r = para.runs[0]
        r.font.size = Pt(16)
        r.font.underline = False  # explicitly NOT underlined

    # ---- Slide 7: Conclusion ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[0])
    slide7.shapes.title.text = 'Thank You'
    slide7.placeholders[1].text = 'Questions & Discussion\nanalytics@company.com | Q4 2024'

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
