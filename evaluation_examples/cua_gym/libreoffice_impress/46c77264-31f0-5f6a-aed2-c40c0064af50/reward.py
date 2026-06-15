"""
Reward Script: Create a combination chart on slide 6 with columns for 'Revenue'
              and a line for 'Profit Margin %' on a secondary axis.
Task ID: impress_tct_060
Domain: libreoffice_impress
Scoring:
  - Component 1: Chart exists on slide 6 (0.20)
  - Component 2: Combination chart with bar + line plots (0.25)
  - Component 3: Revenue series as columns with correct name (0.20)
  - Component 4: Profit Margin % series as line with correct name (0.15)
  - Component 5: Secondary axis (right-side value axis) for line series (0.10)
  - Component 6: Legend present (0.10)
"""

import os
import zipfile
import xml.etree.ElementTree as ET


WORKDIR = '/home/user'
TASK_ID = 'impress_tct_060'


def persist_app_state(domain):
    """Try to save any unsaved state in LibreOffice."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"PRECONDITION FAIL: Need at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # 0-indexed, slide 6

    # Component 1: Chart exists on slide 6 (0.20 points)
    chart_shape = None
    try:
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart_shape = shape
                break

        if chart_shape is not None:
            print(f"PASS: Component 1 — Chart found on slide 6 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — No chart found on slide 6")
            # No chart means no further chart checks can pass
            final_score = min(total_score, 1.0)
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {final_score}")
            return final_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    chart = chart_shape.chart

    # Component 2: Combination chart with both bar and line plots (0.25 points)
    try:
        num_plots = len(chart.plots)
        plot_types = [type(plot).__name__ for plot in chart.plots]
        has_bar = 'BarPlot' in plot_types
        has_line = 'LinePlot' in plot_types

        if num_plots >= 2 and has_bar and has_line:
            print(f"PASS: Component 2 — Combination chart with BarPlot + LinePlot ({plot_types}) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Expected combination chart (BarPlot + LinePlot), found plots: {plot_types}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Revenue series as columns with correct name (0.20 points)
    # We check via XML to get series names reliably
    try:
        ns_c = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        revenue_found = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if n.startswith('ppt/charts/') and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.parse(f).getroot()

                # Look for barChart element containing a series named "Revenue"
                bar_charts = root.findall('.//{%s}barChart' % ns_c)
                for bc in bar_charts:
                    bar_dir = bc.find('{%s}barDir' % ns_c)
                    is_col = bar_dir is not None and bar_dir.get('val') == 'col'
                    series_list = bc.findall('{%s}ser' % ns_c)
                    for ser in series_list:
                        tx = ser.find('{%s}tx' % ns_c)
                        if tx is not None:
                            # Check strCache for the series name
                            str_cache = tx.find('.//{%s}strCache' % ns_c)
                            if str_cache is not None:
                                pts = str_cache.findall('{%s}pt' % ns_c)
                                for pt in pts:
                                    v = pt.find('{%s}v' % ns_c)
                                    if v is not None and v.text and 'revenue' in v.text.lower():
                                        if is_col:
                                            revenue_found = True
                                            print(f"PASS: Component 3 — Revenue series found as column chart (name='{v.text}') (0.20 pts)")

        if revenue_found:
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Revenue series not found as column chart")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Profit Margin % series as line with correct name (0.15 points)
    try:
        margin_found = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if n.startswith('ppt/charts/') and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.parse(f).getroot()

                line_charts = root.findall('.//{%s}lineChart' % ns_c)
                for lc in line_charts:
                    series_list = lc.findall('{%s}ser' % ns_c)
                    for ser in series_list:
                        tx = ser.find('{%s}tx' % ns_c)
                        if tx is not None:
                            str_cache = tx.find('.//{%s}strCache' % ns_c)
                            if str_cache is not None:
                                pts = str_cache.findall('{%s}pt' % ns_c)
                                for pt in pts:
                                    v = pt.find('{%s}v' % ns_c)
                                    if v is not None and v.text and 'profit' in v.text.lower() and 'margin' in v.text.lower():
                                        margin_found = True
                                        print(f"PASS: Component 4 — Profit Margin series found as line chart (name='{v.text}') (0.15 pts)")

        if margin_found:
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Profit Margin % series not found as line chart")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Secondary axis (right-side value axis) for line series (0.10 points)
    try:
        secondary_axis_found = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if n.startswith('ppt/charts/') and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.parse(f).getroot()

                # The line chart should reference axis IDs that correspond to a valAx with axPos="r"
                # Find all valAx elements
                val_axes = root.findall('.//{%s}valAx' % ns_c)
                for vax in val_axes:
                    ax_pos = vax.find('{%s}axPos' % ns_c)
                    delete = vax.find('{%s}delete' % ns_c)
                    is_right = ax_pos is not None and ax_pos.get('val') == 'r'
                    is_visible = delete is None or delete.get('val') == '0'
                    if is_right and is_visible:
                        secondary_axis_found = True
                        print(f"PASS: Component 5 — Secondary value axis found at right position (0.10 pts)")
                        break
                if secondary_axis_found:
                    break

        if secondary_axis_found:
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — No secondary axis (right-side value axis) found")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Legend is present (0.10 points)
    try:
        if chart.has_legend:
            print(f"PASS: Component 6 — Legend present on chart (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — Chart has no legend")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
