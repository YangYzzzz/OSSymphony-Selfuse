"""
Initial Setup: Financial Dashboard presentation with 7 slides.
Slide 6 has title 'Revenue & Profit Margin' but NO chart.
Task ID: impress_tct_060
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_060'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def add_table_slide(prs, title_text, headers, data):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    rows = len(data) + 1
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    table = tbl_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = str(val)

    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    add_title_slide(prs, "Financial Dashboard Q4 2025", "Meridian Technologies Inc. | Confidential")

    # --- Slide 2: Executive Summary ---
    add_content_slide(prs, "Executive Summary", [
        "Total revenue reached $12.8M in Q4 2025, up 14.2% YoY",
        "Gross profit margin improved to 42.3% from 39.1% in Q3",
        "Operating expenses remained flat at $3.9M despite expansion",
        "Net income of $1.52M exceeded forecast by 8.7%",
        "Customer acquisition cost decreased 11% quarter-over-quarter",
        "Three new enterprise contracts signed worth $2.1M annually",
    ])

    # --- Slide 3: Regional Performance (Table) ---
    headers = ["Region", "Revenue ($M)", "Growth (%)", "Market Share (%)", "Active Clients"]
    data = [
        ["North America", "$5.24", "16.8%", "23.4%", "412"],
        ["Europe", "$3.67", "12.1%", "18.7%", "287"],
        ["Asia Pacific", "$2.41", "19.5%", "11.2%", "198"],
        ["Latin America", "$0.89", "8.3%", "6.1%", "74"],
        ["Middle East & Africa", "$0.59", "22.7%", "3.8%", "43"],
    ]
    add_table_slide(prs, "Regional Performance", headers, data)

    # --- Slide 4: Quarterly Trends ---
    add_content_slide(prs, "Quarterly Trends", [
        "Q1 2025: Revenue $10.2M | Margin 38.5% | Headcount 234",
        "Q2 2025: Revenue $11.1M | Margin 39.1% | Headcount 248",
        "Q3 2025: Revenue $11.8M | Margin 40.7% | Headcount 261",
        "Q4 2025: Revenue $12.8M | Margin 42.3% | Headcount 275",
        "Full Year 2025 Total: $45.9M revenue, 40.2% avg margin",
        "Projected Q1 2026: $13.5M revenue target",
    ])

    # --- Slide 5: Expense Breakdown ---
    add_content_slide(prs, "Expense Breakdown", [
        "Salaries & Benefits: $2.14M (54.9% of OpEx)",
        "Cloud Infrastructure: $0.68M (17.4% of OpEx)",
        "Sales & Marketing: $0.52M (13.3% of OpEx)",
        "Office & Facilities: $0.31M (7.9% of OpEx)",
        "R&D Equipment: $0.18M (4.6% of OpEx)",
        "Travel & Entertainment: $0.07M (1.8% of OpEx)",
    ])

    # --- Slide 6: Revenue & Profit Margin (title only, NO chart) ---
    slide6 = add_title_only_slide(prs, "Revenue & Profit Margin")
    # Add a data table as reference for the chart the agent needs to create
    headers6 = ["Quarter", "Revenue ($K)", "Profit Margin (%)"]
    data6 = [
        ["Q1 2024", "2,450", "31.2"],
        ["Q2 2024", "2,780", "33.5"],
        ["Q3 2024", "2,620", "32.1"],
        ["Q4 2024", "3,150", "35.8"],
        ["Q1 2025", "3,380", "37.4"],
        ["Q2 2025", "3,610", "38.9"],
        ["Q3 2025", "3,890", "40.7"],
        ["Q4 2025", "4,230", "42.3"],
    ]
    rows = len(data6) + 1
    cols = len(headers6)
    tbl_shape = slide6.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(5), Inches(4))
    table = tbl_shape.table
    for c, h in enumerate(headers6):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(11)
    for r, row_data in enumerate(data6, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = str(val)

    # --- Slide 7: Outlook & Next Steps ---
    add_content_slide(prs, "Outlook & Next Steps", [
        "Target 15% revenue growth in FY2026",
        "Expand Asia Pacific team by 30 headcount",
        "Launch enterprise tier pricing in Q2 2026",
        "Complete SOC 2 Type II certification by March 2026",
        "Evaluate strategic acquisition opportunities in AI/ML space",
        "Board review meeting scheduled for January 15, 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
