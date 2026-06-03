"""
Initial Setup: Create climate_report.pptx with slide 3 blank (no chart)
Task ID: impress_media_031
Domain: libreoffice_impress
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'impress_media_031'
# Agent-facing file (what the task instruction references)
AGENT_FILE = f'{DESKTOP}/climate_report.pptx'
# Initial file for reward-gen tracking
OUTPUT = f'{WORKDIR}/{TASK_ID}_initial.pptx'

def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide layout
    slide1.shapes.title.text = "Climate Report H1 2024"
    slide1.placeholders[1].text = "A comprehensive analysis of temperature trends\nacross major Chinese cities"

    # ---- Slide 2: Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide2.shapes.title.text = "Overview"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "This report covers temperature data from January to June 2024."
    p2a = tf2.add_paragraph()
    p2a.text = "Cities analyzed: Beijing, Shanghai, Guangzhou"
    p2a.level = 1
    p2b = tf2.add_paragraph()
    p2b.text = "Data sourced from National Meteorological Administration"
    p2b.level = 1
    p2c = tf2.add_paragraph()
    p2c.text = "Average monthly temperatures recorded in degrees Celsius"
    p2c.level = 1

    # ---- Slide 3: Climate Data (blank except for title text box, NO chart) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box only — no chart present (agent must add it)
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = title_box.text_frame
    tf3.text = "Climate Data"
    p_title = tf3.paragraphs[0]
    run_title = p_title.runs[0]
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # ---- Slide 4: Summary ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide4.shapes.title.text = "Summary & Key Findings"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Beijing experienced the highest temperature variation (range: 3\u00b0C to 29\u00b0C)"
    p4a = tf4.add_paragraph()
    p4a.text = "Guangzhou remained the warmest city throughout H1 2024"
    p4a.level = 1
    p4b = tf4.add_paragraph()
    p4b.text = "Shanghai showed moderate and stable temperature increases"
    p4b.level = 1
    p4c = tf4.add_paragraph()
    p4c.text = "All three cities show consistent warming trend from Jan to Jun"
    p4c.level = 1

    # ---- Slide 5: Conclusions ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide5.shapes.title.text = "Conclusions"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Urban heat island effect observed in all cities"
    p5a = tf5.add_paragraph()
    p5a.text = "Seasonal patterns align with historical records"
    p5a.level = 1
    p5b = tf5.add_paragraph()
    p5b.text = "Further monitoring recommended for Q3-Q4 2024"
    p5b.level = 1

    # Ensure Desktop directory exists
    os.makedirs(DESKTOP, exist_ok=True)

    # Save as agent-facing file
    prs.save(AGENT_FILE)
    print(f'Agent file created: {AGENT_FILE}')

    # Also save as initial tracking file for reward-gen
    shutil.copy(AGENT_FILE, OUTPUT)
    print(f'Initial tracking file created: {OUTPUT}')
    print(f'Slides: {len(prs.slides)}')
    print(f'Slide 3 shape count (should be 1 text box, no chart): {len(prs.slides[2].shapes)}')

create_initial()
