"""
Reward Script: Insert a line chart on slide 3 showing temperature trends for three cities over 6 months,
               then add a title 'Temperature Trends H1 2024' and show the legend.
Task ID: impress_media_031
Domain: libreoffice_impress
Scoring:
  Component 1: Chart exists on slide 3 (0.30 pts)
  Component 2: Chart type is Lines Only (LINE) and chart title is 'Temperature Trends H1 2024' (0.30 pts)
  Component 3: Chart has 3 series (Beijing, Shanghai, Guangzhou) with 6 data points each (0.25 pts)
  Component 4: Legend is visible/shown (0.15 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

WORKDIR = '/home/user'  # VM path — all reward scripts run on the VM
TASK_ID = 'impress_media_031'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3

    # Component 1: A chart shape exists on slide 3 (0.30 points)
    # This FAILS on initial (no chart) -> PASSES on golden (chart present)
    chart_shape = None
    try:
        for shape in slide3.shapes:
            if shape.shape_type == 3:  # MSO_SHAPE_TYPE.CHART == 3
                chart_shape = shape
                break
        if chart_shape is not None:
            print(f"PASS: Component 1 — Chart shape found on slide 3 (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — No chart shape found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no chart found, remaining components cannot be evaluated
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart type is LINE ('Lines Only') AND chart title is correct (0.30 points)
    # Both fail on initial (no chart at all) -> pass on golden
    try:
        chart_type_ok = (chart.chart_type == XL_CHART_TYPE.LINE)
        title_ok = False
        if chart.has_title:
            actual_title = chart.chart_title.text_frame.text.strip()
            title_ok = (actual_title == 'Temperature Trends H1 2024')
            if not title_ok:
                print(f"FAIL: Component 2 — Chart title expected 'Temperature Trends H1 2024', found {actual_title!r}")
        else:
            print(f"FAIL: Component 2 — Chart has no title")

        if not chart_type_ok:
            print(f"FAIL: Component 2 — Chart type expected LINE (4), found {chart.chart_type!r}")

        if chart_type_ok and title_ok:
            print(f"PASS: Component 2 — Chart type is LINE and title is 'Temperature Trends H1 2024' (0.30 pts)")
            total_score += 0.30
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Chart has 3 series (Beijing, Shanghai, Guangzhou) each with 6 data points (0.25 points)
    # Fails on initial (no chart) -> passes on golden
    try:
        expected_series = ['Beijing', 'Shanghai', 'Guangzhou']
        series_list = list(chart.series)
        series_count_ok = (len(series_list) == 3)
        series_names_ok = False
        series_points_ok = False

        if series_count_ok:
            actual_names = [s.name for s in series_list]
            series_names_ok = (actual_names == expected_series)
            if not series_names_ok:
                print(f"FAIL: Component 3 — Series names expected {expected_series}, found {actual_names}")
            # Check each series has 6 data points
            points_check = []
            for s in series_list:
                vals = list(s.values)
                points_check.append(len(vals) == 6)
            series_points_ok = all(points_check)
            if not series_points_ok:
                counts = [len(list(s.values)) for s in series_list]
                print(f"FAIL: Component 3 — Expected 6 data points per series, found: {counts}")
        else:
            print(f"FAIL: Component 3 — Expected 3 series, found {len(series_list)}")

        if series_count_ok and series_names_ok and series_points_ok:
            print(f"PASS: Component 3 — 3 series (Beijing, Shanghai, Guangzhou) each with 6 data points (0.25 pts)")
            total_score += 0.25
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Legend is visible (0.15 points)
    # Fails on initial (no chart) -> passes on golden (chart has legend)
    try:
        legend_ok = chart.has_legend
        if legend_ok:
            print(f"PASS: Component 4 — Legend is visible (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — Legend is not shown (chart.has_legend == False)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against golden file (path on VM)
file_path = f'{WORKDIR}/{TASK_ID}_initial.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
