"""
Reward Script: Insert an embedded line chart on slide 4 with specific data, title, gridlines, markers, and data label.
Task ID: impress_gf2_005
Domain: libreoffice_impress
Scoring:
  Component 1: Line chart exists on slide 4 (0.25)
  Component 2: Chart data matches (6 categories Jan-Jun, correct values) (0.25)
  Component 3: Chart title is 'H1 Revenue Trend' (0.15)
  Component 4: Data point markers are visible (0.10)
  Component 5: Major gridlines present (0.10)
  Component 6: Data label '134K' on last point (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_005'

EXPECTED_CATEGORIES = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']
EXPECTED_VALUES = [85000.0, 92000.0, 78000.0, 105000.0, 118000.0, 134000.0]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # slide 4 (0-indexed)

    # Find chart shape on slide 4
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Line chart exists on slide 4 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            chart_type = chart.chart_type
            # Accept LINE (64), LINE_MARKERS (65), LINE_STACKED (67), LINE_MARKERS_STACKED (68)
            is_line = chart_type in (64, 65, 67, 68)
            if is_line:
                print(f"PASS: Component 1 -- Line chart found on slide 4, type={chart_type} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Chart on slide 4 is not a line chart, type={chart_type}")
        else:
            print("FAIL: Component 1 -- No chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chart, remaining components cannot pass
    if chart_shape is None:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    chart = chart_shape.chart

    # Component 2: Chart data matches (0.25 points)
    try:
        plot = chart.plots[0]
        series = plot.series[0]
        actual_values = list(series.values)
        actual_categories = [str(c) for c in plot.categories]

        cats_match = actual_categories == EXPECTED_CATEGORIES
        vals_match = len(actual_values) == len(EXPECTED_VALUES)
        if vals_match:
            for av, ev in zip(actual_values, EXPECTED_VALUES):
                if abs(av - ev) > 1.0:
                    vals_match = False
                    break

        if cats_match and vals_match:
            print(f"PASS: Component 2 -- Chart data matches: categories={actual_categories}, values={actual_values} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 -- Data mismatch. Categories: {actual_categories} (expected {EXPECTED_CATEGORIES}), Values: {actual_values} (expected {EXPECTED_VALUES})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart title is 'H1 Revenue Trend' (0.15 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            if title_text == 'H1 Revenue Trend':
                print(f"PASS: Component 3 -- Chart title is 'H1 Revenue Trend' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 -- Chart title is '{title_text}', expected 'H1 Revenue Trend'")
        else:
            print("FAIL: Component 3 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Data point markers are visible (0.10 points)
    try:
        plot = chart.plots[0]
        series = plot.series[0]
        marker = series.marker
        # Marker style != None and != XL_MARKER_STYLE.NONE (which is -4142)
        marker_style = marker.style
        if marker_style is not None and marker_style != -4142:
            print(f"PASS: Component 4 -- Markers visible, style={marker_style} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 -- No markers visible, style={marker_style}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Major gridlines present on value axis (0.10 points)
    try:
        ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart'}
        has_gridlines = False
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if n.startswith('ppt/charts/') and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.parse(f).getroot()
                    gridlines = root.findall('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/chart}majorGridlines')
                    if gridlines:
                        has_gridlines = True
                        break

        if has_gridlines:
            print(f"PASS: Component 5 -- Major gridlines present (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 5 -- No major gridlines found in chart XML")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Data label '134K' on last data point (index 5) (0.15 points)
    try:
        ns_c = '{http://schemas.openxmlformats.org/drawingml/2006/chart}'
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        found_label = False

        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_files = [n for n in zf.namelist() if n.startswith('ppt/charts/') and n.endswith('.xml')]
            for cf in chart_files:
                with zf.open(cf) as f:
                    root = ET.parse(f).getroot()
                    # Search for dLbl elements (individual data labels)
                    for dlbl in root.iter(f'{ns_c}dLbl'):
                        idx_el = dlbl.find(f'{ns_c}idx')
                        if idx_el is not None and idx_el.get('val') == '5':
                            # Found label on index 5 (June = last point)
                            # Check if it contains '134K' text
                            label_text = ''
                            for t_el in dlbl.iter(f'{ns_a}t'):
                                if t_el.text:
                                    label_text += t_el.text
                            if '134K' in label_text or '134k' in label_text.lower():
                                found_label = True
                                break
                    # Also check dLbls (collection-level data labels) with showVal
                    if not found_label:
                        for dlbls in root.iter(f'{ns_c}dLbls'):
                            show_val = dlbls.find(f'{ns_c}showVal')
                            if show_val is not None and show_val.get('val') == '1':
                                # All values shown; check if the last point shows 134000 or 134K
                                # This is acceptable if values are displayed for all points
                                # But we specifically need '134K' text, so only pass on explicit label
                                pass

        if found_label:
            print(f"PASS: Component 6 -- Data label '134K' found on last point (index 5) (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 6 -- Data label '134K' not found on last data point (index 5)")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
