"""
Initial Setup: Create H1 Financial presentation with 8 slides.
Slide 4 has title 'Revenue Trend Analysis' and empty content area (no chart).
Task ID: impress_gf2_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bullet body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "H1 Financial Report"
    slide1.placeholders[1].text = "Fiscal Year 2025 — First Half Review"

    # --- Slide 2: Executive Summary ---
    add_title_body_slide(prs, "Executive Summary", [
        "Total H1 revenue reached $612K, a 14% increase over prior year",
        "Operating margin improved to 23.5%, up from 19.8%",
        "Customer acquisition cost decreased by 11% through organic growth",
        "New enterprise contracts signed with 3 Fortune 500 companies",
        "Employee headcount grew from 142 to 178 across all departments",
    ])

    # --- Slide 3: Expense Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Title text box
    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Expense Breakdown by Department"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Table
    rows, cols = 7, 4
    tbl = slide3.shapes.add_table(rows, cols, Inches(1.5), Inches(1.6), Inches(10), Inches(4.2))
    table = tbl.table
    headers = ["Department", "Q1 Spend ($K)", "Q2 Spend ($K)", "Change (%)"]
    data = [
        ["Engineering", "186", "198", "+6.5%"],
        ["Marketing", "94", "112", "+19.1%"],
        ["Sales", "78", "85", "+9.0%"],
        ["Operations", "52", "48", "-7.7%"],
        ["HR & Admin", "31", "34", "+9.7%"],
        ["R&D", "67", "73", "+9.0%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        solidFill.append(solidFill.makeelement(qn('a:srgbClr'), {'val': '1F3864'}))
        tcPr.append(solidFill)

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 4: Revenue Trend Analysis (EMPTY - no chart) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title
    title_box = slide4.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.9))
    tf4 = title_box.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Revenue Trend Analysis"
    p4.font.size = Pt(32)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Placeholder text indicating empty content area
    content_box = slide4.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3.5))
    tf_content = content_box.text_frame
    p_content = tf_content.paragraphs[0]
    p_content.text = "[Chart to be inserted here]"
    p_content.alignment = PP_ALIGN.CENTER
    p_content.font.size = Pt(18)
    p_content.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)
    p_content.font.italic = True

    # --- Slide 5: Q2 Projections ---
    add_title_body_slide(prs, "Q2 Projections", [
        "Projected Q3 revenue: $155K based on current pipeline",
        "New product launch expected to contribute $28K monthly by August",
        "Customer retention rate holding steady at 94.2%",
        "Hiring plan: 12 new positions across Engineering and Sales",
        "Capital expenditure budget approved at $340K for H2",
    ])

    # --- Slide 6: Regional Performance ---
    add_title_body_slide(prs, "Regional Performance", [
        "North America: $387K (63.2% of total), +18% YoY",
        "Europe: $128K (20.9%), +9% YoY",
        "Asia Pacific: $72K (11.8%), +22% YoY — fastest growing region",
        "Latin America: $25K (4.1%), +5% YoY — early stage market",
        "EMEA expansion on track with Munich office opening in Q3",
    ])

    # --- Slide 7: Key Takeaways ---
    add_title_body_slide(prs, "Key Takeaways", [
        "Strong revenue growth trajectory with accelerating monthly gains",
        "Cost optimization efforts showing measurable results across departments",
        "Geographic diversification reducing single-market risk",
        "Product roadmap aligned with enterprise customer feedback",
        "Board-approved targets for H2: $720K revenue, 25% operating margin",
    ])

    # --- Slide 8: Thank You ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Discussion"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
