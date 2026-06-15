"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’ve got a 40-slide Impress deck and need every single text box—titles, bullets, footers, the works—to use the exact font “Liberation Sans Narrow”. What’s the quickest way to apply that across the board via the master slide or styles so I don’t have to touch each slide individually?
Generated: 2025-09-10 12:31:04
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import zipfile
from pptx import Presentation
from lxml import etree

"""Reward script for verifying that every text box in a 40-slide (or larger) PPTX
presentation uses the exact font “Liberation Sans Narrow”.

Scoring (progressive):
1. Slide-count requirement (>=40) .................................... 0.1 points
2. Font correctness on every text run ............................... up to 0.9 points
   • The 0.9 points are linearly scaled by the proportion of text runs
     that use the required font. A text run is considered correct if:
       a) It explicitly specifies the expected font, **or**
       b) It relies on the slide theme **and** the theme’s major & minor
          latin fonts are the expected font.
The script prints detailed diagnostics and the final score as
“REWARD: X.X” (float between 0.0 and 1.0).
"""

EXPECTED_FONT = "Liberation Sans Narrow"
MIN_SLIDES    = 40

# XML namespaces used inside PPTX files
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}

def _get_theme_fonts(zip_obj):
    """Return (major_font, minor_font) latin typefaces from the first theme file."""
    theme_path = next((n for n in zip_obj.namelist()
                       if n.startswith('ppt/theme/') and n.endswith('.xml')), None)
    if not theme_path:
        return None, None

    root = etree.fromstring(zip_obj.read(theme_path))
    major_el = root.xpath('.//a:fontScheme/a:majorFont/a:latin', namespaces=NS)
    minor_el = root.xpath('.//a:fontScheme/a:minorFont/a:latin', namespaces=NS)
    major_font = major_el[0].get('typeface') if major_el else None
    minor_font = minor_el[0].get('typeface') if minor_el else None
    return major_font, minor_font


def _iter_text_runs(zip_obj):
    """Yield tuples (slide_path, explicit_font_or_None, uses_theme_bool)."""
    slide_files = [n for n in zip_obj.namelist()
                   if n.startswith('ppt/slides/slide') and n.endswith('.xml')]

    for slide_file in slide_files:
        root = etree.fromstring(zip_obj.read(slide_file))
        for run in root.xpath('.//a:r', namespaces=NS):
            rPr = run.find('.//a:rPr', namespaces=NS)
            if rPr is None:
                # No run properties: inherits theme
                yield slide_file, None, True
                continue

            latin = rPr.find('.//a:latin', namespaces=NS)
            if latin is not None and latin.get('typeface'):
                tf = latin.get('typeface')
                # “+mn-lt”, “+mj-lt”, etc. are placeholders meaning “use theme font”
                if tf.startswith('+'):
                    yield slide_file, None, True
                else:
                    yield slide_file, tf, False
            else:
                # No latin child → inherits theme
                yield slide_file, None, True


def verify_presentation_fonts(file_path: str,
                              expected_font: str = EXPECTED_FONT,
                              min_slides: int = MIN_SLIDES) -> float:
    """Main verification function. Returns a float score between 0.0 and 1.0."""

    print(f"Verifying PPTX font usage for: {file_path}")
    print(f"Expected font: {expected_font}\n")

    # --- Preliminary checks -------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)  # validates file – no points awarded
        slide_count = len(prs.slides)
        print(f"Loaded presentation with {slide_count} slides")
    except Exception as e:
        print(f"✗ Cannot open PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- Scoring initialisation --------------------------------------------
    total_score = 0.0

    # 1) Slide count requirement (0.1 points)
    if slide_count >= min_slides:
        total_score += 0.1
        print(f"✓ Slide-count requirement met (≥{min_slides}) (+0.1)")
    else:
        print(f"✗ Slide-count requirement NOT met (found {slide_count} < {min_slides})")

    # --- Deep XML analysis --------------------------------------------------
    with zipfile.ZipFile(file_path, 'r') as zf:
        # Theme fonts
        major_font, minor_font = _get_theme_fonts(zf)
        theme_ok = (major_font and minor_font and
                    major_font.lower() == expected_font.lower() and
                    minor_font.lower() == expected_font.lower())

        if theme_ok:
            print(f"✓ Theme major & minor latin fonts set to '{expected_font}'")
        else:
            print(f"Theme fonts not fully set to expected font.")
            print(f"  Major latin: {major_font}")
            print(f"  Minor latin: {minor_font}")

        # Iterate through every text run on every slide
        total_runs   = 0
        correct_runs = 0
        for slide_path, explicit_font, uses_theme in _iter_text_runs(zf):
            total_runs += 1
            if uses_theme:
                # Correct only if the theme itself is correct
                if theme_ok:
                    correct_runs += 1
            else:
                if explicit_font and explicit_font.lower() == expected_font.lower():
                    correct_runs += 1

    if total_runs == 0:
        print("✗ No text runs detected – cannot evaluate font usage")
        print("REWARD: 0.0")
        return 0.0

    proportion_correct = correct_runs / total_runs
    font_score         = round(0.9 * proportion_correct, 2)  # up to 0.9 pts
    total_score       += font_score

    # --- Diagnostics --------------------------------------------------------
    print(f"Total text runs analysed: {total_runs}")
    print(f"Runs using correct font: {correct_runs}")
    print(f"Font correctness proportion: {proportion_correct:.2%} (+{font_score})")

    # --- Final score capping & output --------------------------------------
    final_score = round(min(1.0, total_score), 2)
    print(f"\nFINAL SCORE: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# Execute verification when the script is run directly
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/ive_got_a_40_slide_impress_deck_and_need_every_single_text_boxtitles_bullets_footers_the_worksto_use_golden.pptx"
    verify_presentation_fonts(FILE_PATH)

