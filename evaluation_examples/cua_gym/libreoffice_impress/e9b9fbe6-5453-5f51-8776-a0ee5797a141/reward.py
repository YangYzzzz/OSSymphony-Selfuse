"""
FINAL REWARD SCRIPT - SUCCESS
Task: Hey there! I need to include our motto, 'Innovation First,' at the bottom of all slides in our presentation, but I want to keep the title slide clean without it. Can you guide me through the steps to do this in LibreOffice Impress?
Generated: 2025-08-07 08:39:55
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_motto_footer(file_path, motto):
    print(f"Checking presentation file: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File existence and load (0.2 points)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score:.3f}")
        return
    try:
        prs = Presentation(file_path)
        num_slides = len(prs.slides)
        print(f"✓ File loaded successfully. Slide count: {num_slides} (0.2)")
        total_score += 0.2
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: {total_score:.3f}")
        return

    expected_text = motto.strip().lower()

    # Requirement 2: Motto absent on title slide (0.1 points)
    title_slide = prs.slides[0]
    found_in_title = False
    for shape in title_slide.shapes:
        if hasattr(shape, "text") and expected_text in shape.text.strip().lower():
            found_in_title = True
            print(f"✗ Motto found on title slide: '{shape.text.strip()}' (0 points for this requirement)")
            break
    if not found_in_title:
        print(f"✓ Motto absent on title slide (0.1)")
        total_score += 0.1

    # Requirement 3: Motto present on all other slides (0.7 points total)
    slides_to_check = num_slides - 1
    if slides_to_check > 0:
        slide_weight = 0.7 / slides_to_check
        for idx in range(1, num_slides):
            slide = prs.slides[idx]
            motto_found = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and expected_text in shape.text.strip().lower():
                    motto_found = True
                    print(f"✓ Motto found on slide {idx+1} (+{slide_weight:.3f})")
                    total_score += slide_weight
                    break
            if not motto_found:
                print(f"✗ Motto missing on slide {idx+1} (0 points)")
    else:
        # No non-title slides: assign full points for presence
        print("✓ No non-title slides to check; assigning full points for motto presence (0.7)")
        total_score += 0.7

    # Final scoring and output
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score:.3f} / {max_score}")
    print(f"REWARD: {final_score:.3f}")

if __name__ == '__main__':
    file_path = '/home/user/hey_there_i_need_to_include_our_motto_innovation_first_at_the_bottom_of_all_slides_in_our_presentati.pptx'
    verify_motto_footer(file_path, 'Innovation First')
