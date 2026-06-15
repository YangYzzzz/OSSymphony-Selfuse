"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 257 feels way too cramped. In LibreOffice Impress, how do I make the text in the content box use EXACTLY 22 pt line spacing and then add 10 pt of space after every paragraph so it breathes a bit?
Generated: 2025-09-10 19:00:11
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER


def _collect_content_paragraphs(slide):
    """Return a list of paragraphs from the non-title text shapes on the slide."""
    paragraphs = []
    for shape in slide.shapes:
        # Only shapes that actually have a text_frame are relevant
        if not hasattr(shape, "text_frame"):
            continue

        # Skip title / subtitle placeholders – task targets the content box
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
                PP_PLACEHOLDER.SUBTITLE,
            ):
                continue

        for para in shape.text_frame.paragraphs:
            # Ignore empty paragraphs – LibreOffice often keeps a trailing empty one
            if para.text and para.text.strip():
                paragraphs.append(para)
    return paragraphs


def _score_paragraph_formatting(paragraphs, expected_line, expected_space, tolerance):
    """Return (line_score, space_score) each in the range [0,1]."""
    if not paragraphs:
        return 0.0, 0.0

    correct_lines = 0
    correct_spaces = 0
    total = len(paragraphs)

    for para in paragraphs:
        line_val = para.line_spacing
        space_val = para.space_after

        if line_val is not None and abs(line_val - expected_line) <= tolerance:
            correct_lines += 1
        if space_val is not None and abs(space_val - expected_space) <= tolerance:
            correct_spaces += 1

    line_score = correct_lines / total
    space_score = correct_spaces / total
    return line_score, space_score


def verify_slide_257_formatting(file_path):
    """Verify that slide 257 content paragraphs use 22 pt line spacing and 10 pt space-after.

    Progressive scoring:
      – Up to 0.5 for correct line spacing proportionally
      – Up to 0.5 for correct space-after proportionally
    Returns a float in [0,1] and prints detailed diagnostics plus
    a final line in the form `REWARD: X.X` required by the evaluation harness.
    """
    print(f"Verifying formatting for slide 257 in: {file_path}")

    # 1. Preliminary checks --------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File does not exist.")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to open presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    target_index = 256  # zero-based index for slide 257
    if len(prs.slides) <= target_index:
        print(f"✗ Presentation has only {len(prs.slides)} slides; slide 257 missing.")
        print("REWARD: 0.0")
        return 0.0

    # 2. Parameter setup -----------------------------------------------------
    expected_line = int(Pt(22))   # centipoints (1/100 of a point)
    expected_space = int(Pt(10))  # centipoints
    tolerance = 1270              # ±0.1 pt tolerance => 127 centipoints; use 1270 (±1 pt) for safety

    # 3. Collect paragraphs to evaluate -------------------------------------
    slide = prs.slides[target_index]
    paragraphs = _collect_content_paragraphs(slide)

    if not paragraphs:
        print("✗ No text paragraphs found in content box of slide 257.")
        print("REWARD: 0.0")
        return 0.0

    # 4. Score formatting ----------------------------------------------------
    line_score, space_score = _score_paragraph_formatting(
        paragraphs, expected_line, expected_space, tolerance
    )

    # Detailed diagnostics for transparency
    print(f"Paragraphs evaluated: {len(paragraphs)}")
    print(f"Line spacing correctness: {line_score*100:.1f}%")
    print(f"Space-after correctness: {space_score*100:.1f}%")

    # 5. Aggregate final score ----------------------------------------------
    final_score = round(min(0.5 * line_score + 0.5 * space_score, 1.0), 2)
    print(f"REWARD: {final_score}")
    return final_score


# ---------------------------------------------------------------------------
# Execute verification when run as a script
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/slide_257_feels_way_too_cramped_in_libreoffice_impress_"
        "how_do_i_make_the_text_in_the_content_box_use_golden.pptx"
    )
    verify_slide_257_formatting(FILE_PATH)

