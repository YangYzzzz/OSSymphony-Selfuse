"""
Initial Setup: Create a 10-slide Research Report presentation with menu text boxes on slide 2.
Task ID: impress_gf2_003
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title(slide, text):
    """Add a title text box to the top of a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    return txBox


def add_body(slide, text):
    """Add a body text box below the title area."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout index - find it
    blank_layout = prs.slide_layouts[6]  # Usually blank in default template

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(blank_layout)
    txBox = slide1.shapes.add_textbox(Inches(2), Inches(1.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Research Report"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    txBox2 = slide1.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Annual Findings & Analysis\nDr. Sarah Chen | March 2025"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0x59, 0x56, 0x59)

    # ---- Slide 2: Menu Slide (4 text boxes, NO hyperlinks) ----
    slide2 = prs.slides.add_slide(blank_layout)
    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Table of Contents"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    menu_items = ['Introduction', 'Methodology', 'Results', 'Conclusion']
    box_width = Inches(5)
    box_height = Inches(0.8)
    start_x = (prs.slide_width - box_width) // 2
    start_y = Inches(2.0)
    spacing = Inches(1.2)

    for i, label in enumerate(menu_items):
        txBox = slide2.shapes.add_textbox(
            start_x,
            start_y + Emu(int(spacing * i)),
            box_width,
            box_height,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xD6, 0xE4, 0xF0)

    # ---- Slide 3: Introduction ----
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "Introduction")
    add_body(slide3, (
        "This research examines the impact of renewable energy adoption on "
        "urban air quality metrics across 15 metropolitan areas from 2020 to 2024.\n\n"
        "Key objectives include quantifying particulate matter reduction, "
        "assessing economic co-benefits, and identifying policy drivers that "
        "accelerate clean energy transitions in densely populated regions."
    ))

    # ---- Slide 4: Background ----
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "Background")
    add_body(slide4, (
        "Previous studies have shown a strong correlation between fossil fuel "
        "dependency and elevated PM2.5 concentrations in urban environments.\n\n"
        "The WHO estimates that ambient air pollution contributes to 4.2 million "
        "premature deaths annually, with transportation and power generation "
        "accounting for the largest emission shares."
    ))

    # ---- Slide 5: Methodology ----
    slide5 = prs.slides.add_slide(blank_layout)
    add_title(slide5, "Methodology")
    add_body(slide5, (
        "Data Collection:\n"
        "- EPA Air Quality System (AQS) monitoring data\n"
        "- EIA renewable energy capacity reports\n"
        "- Census Bureau population density statistics\n\n"
        "Analysis Framework:\n"
        "- Panel regression with fixed effects\n"
        "- Difference-in-differences for policy intervention analysis\n"
        "- Spatial autocorrelation testing (Moran's I)"
    ))

    # ---- Slide 6: Data Overview ----
    slide6 = prs.slides.add_slide(blank_layout)
    add_title(slide6, "Data Overview")
    add_body(slide6, (
        "Sample: 15 metropolitan areas, quarterly observations (2020 Q1 - 2024 Q4)\n"
        "Total observations: 300 city-quarter data points\n\n"
        "Variables measured:\n"
        "- PM2.5 concentration (micrograms per cubic meter)\n"
        "- Renewable energy share (%)\n"
        "- GDP per capita\n"
        "- Vehicle registrations per 1000 residents"
    ))

    # ---- Slide 7: Results ----
    slide7 = prs.slides.add_slide(blank_layout)
    add_title(slide7, "Results")
    add_body(slide7, (
        "Key Findings:\n"
        "- A 10% increase in renewable energy share correlates with a 3.2% "
        "reduction in PM2.5 (p < 0.01)\n"
        "- Cities with carbon pricing showed 40% faster adoption rates\n"
        "- Economic co-benefits averaged $1.8B annually per metro area\n\n"
        "The strongest effects were observed in cities that combined renewable "
        "mandates with electric vehicle incentive programs."
    ))

    # ---- Slide 8: Discussion ----
    slide8 = prs.slides.add_slide(blank_layout)
    add_title(slide8, "Discussion")
    add_body(slide8, (
        "Our findings align with the broader literature on energy transition benefits, "
        "while providing novel evidence on the synergistic effects of multi-policy "
        "approaches.\n\n"
        "Limitations include potential endogeneity in policy adoption and the "
        "relatively short time horizon. Future work should extend the analysis "
        "to 2030 targets and include indoor air quality metrics."
    ))

    # ---- Slide 9: Conclusion ----
    slide9 = prs.slides.add_slide(blank_layout)
    add_title(slide9, "Conclusion")
    add_body(slide9, (
        "Renewable energy adoption significantly improves urban air quality.\n\n"
        "Policy Recommendations:\n"
        "- Implement combined renewable + EV incentive frameworks\n"
        "- Establish carbon pricing at $45-65 per ton\n"
        "- Invest in grid modernization for distributed energy resources\n"
        "- Set PM2.5 reduction targets tied to renewable portfolio standards"
    ))

    # ---- Slide 10: References ----
    slide10 = prs.slides.add_slide(blank_layout)
    add_title(slide10, "References")
    add_body(slide10, (
        "1. WHO (2023). Ambient Air Pollution: A Global Assessment.\n"
        "2. EPA (2024). Air Quality System Annual Summary Report.\n"
        "3. EIA (2024). Monthly Energy Review, Table 7.2a.\n"
        "4. Myung, J. et al. (2023). Renewable Energy and Urban Health. Nature Energy, 8(4), 312-325.\n"
        "5. Rodriguez, P. & Kim, S. (2022). Carbon Pricing Effectiveness in Metropolitan Areas. JAERE, 9(2), 88-117."
    ))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
