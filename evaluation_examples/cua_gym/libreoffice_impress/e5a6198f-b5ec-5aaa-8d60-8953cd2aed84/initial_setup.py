"""
Initial Setup: Create a 5-slide presentation with slide 4 having a bulleted list (default bullet size).
Task ID: impstruct_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Review"
    slide1.placeholders[1].text = "Prepared by the Strategy & Operations Team"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Market Overview & Competitive Landscape"
    items_s2 = [
        "Revenue Performance by Region",
        "Product Roadmap Update",
        "Key Points for Leadership",
        "Next Steps & Action Items",
    ]
    for item in items_s2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Market Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Global SaaS market grew 18.4% YoY reaching $232B in Q2 2025"
    items_s3 = [
        "North America leads with 42% market share",
        "APAC fastest growth region at 24.7% YoY",
        "Enterprise segment accounts for 61% of new bookings",
        "Mid-market showing strong momentum with 15% growth",
        "SMB segment stabilizing after Q1 softness",
    ]
    for item in items_s3:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Key Points (the target slide with bullets) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Points"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Accelerate cloud migration timeline from Q4 to Q3 2025"
    items_s4 = [
        "Expand partner ecosystem with 3 new strategic alliances",
        "Invest $4.2M in AI-driven customer engagement platform",
        "Consolidate EMEA operations under unified regional structure",
    ]
    for item in items_s4:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Timeline"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Complete vendor evaluation by June 15, 2025"
    items_s5 = [
        "Schedule executive alignment session for week of June 23",
        "Finalize regional restructuring proposal by July 1",
        "Launch pilot program with selected enterprise accounts",
        "Report Q3 progress at September board meeting",
    ]
    for item in items_s5:
        p = body5.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
