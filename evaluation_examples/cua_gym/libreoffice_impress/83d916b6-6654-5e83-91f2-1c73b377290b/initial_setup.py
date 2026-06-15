"""
Initial Setup: Create a 5-slide presentation with a car icon on slide 4 (no animations).
Task ID: impress_ma_077
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_077'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_car_shape(slide, left, top, width, height):
    """
    Add a car-like icon shape to a slide using a combination of shapes
    grouped visually. We'll use a simple rectangle + two circles (wheels)
    to represent a car icon.
    """
    # Car body - a rounded rectangle
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, int(height * 0.6)
    )
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(0xE8, 0x3E, 0x3E)  # Red car
    body.line.fill.background()
    body.name = "Car Icon"

    # Car roof - smaller rectangle on top
    roof_w = int(width * 0.5)
    roof_h = int(height * 0.35)
    roof_left = left + int(width * 0.2)
    roof_top = top - int(height * 0.25)
    roof = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        roof_left, roof_top, roof_w, roof_h
    )
    roof.fill.solid()
    roof.fill.fore_color.rgb = RGBColor(0xC0, 0x30, 0x30)  # Darker red
    roof.line.fill.background()
    roof.name = "Car Roof"

    # Left wheel
    wheel_r = int(height * 0.25)
    lw = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + int(width * 0.15) - wheel_r // 2,
        top + int(height * 0.5),
        wheel_r, wheel_r
    )
    lw.fill.solid()
    lw.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    lw.line.fill.background()
    lw.name = "Left Wheel"

    # Right wheel
    rw = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + int(width * 0.75) - wheel_r // 2,
        top + int(height * 0.5),
        wheel_r, wheel_r
    )
    rw.fill.solid()
    rw.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    rw.line.fill.background()
    rw.name = "Right Wheel"

    # Windshield
    ws_w = int(width * 0.2)
    ws_h = int(height * 0.25)
    ws = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        roof_left + int(roof_w * 0.6),
        roof_top + int(roof_h * 0.15),
        ws_w, ws_h
    )
    ws.fill.solid()
    ws.fill.fore_color.rgb = RGBColor(0xAD, 0xD8, 0xE6)  # Light blue glass
    ws.line.fill.background()
    ws.name = "Windshield"

    return body  # Return the main car body shape for animation targeting


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Fleet Management Overview"
    slide1.placeholders[1].text = "Q1 2025 Vehicle Operations Report"

    # --- Slide 2: Vehicle Fleet Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Vehicle Fleet Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Current Fleet Composition"
    items = [
        "Sedans: 42 units (35% of fleet)",
        "SUVs: 28 units (23% of fleet)",
        "Delivery Vans: 31 units (26% of fleet)",
        "Electric Vehicles: 19 units (16% of fleet)",
        "Total: 120 active vehicles across 5 regional depots",
    ]
    for item in items:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 1
        p.space_after = Pt(4)

    # --- Slide 3: Performance Metrics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    title3.text_frame.paragraphs[0].text = "Q1 Performance Metrics"
    title3.text_frame.paragraphs[0].font.size = Pt(28)
    title3.text_frame.paragraphs[0].font.bold = True

    # Add a simple metrics table
    tbl_shape = slide3.shapes.add_table(5, 3, Inches(1), Inches(1.5), Inches(7), Inches(3))
    tbl = tbl_shape.table
    headers = ["Metric", "Target", "Actual"]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    data = [
        ["Average Fuel Efficiency", "28.5 mpg", "29.1 mpg"],
        ["On-Time Delivery Rate", "95%", "97.2%"],
        ["Maintenance Compliance", "100%", "98.5%"],
        ["Vehicle Utilization", "85%", "88.3%"],
    ]
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val

    # --- Slide 4: Car Animation Demo (car icon at left, NO animations) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    title4.text_frame.paragraphs[0].text = "Route Simulation"
    title4.text_frame.paragraphs[0].font.size = Pt(28)
    title4.text_frame.paragraphs[0].font.bold = True

    # Subtitle
    sub4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(8), Inches(0.5))
    sub4.text_frame.paragraphs[0].text = "Click to animate the vehicle along its delivery route"
    sub4.text_frame.paragraphs[0].font.size = Pt(16)
    sub4.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Road line (a thin rectangle across the slide)
    road = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.0), Inches(9.0), Inches(0.08)
    )
    road.fill.solid()
    road.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    road.line.fill.background()
    road.name = "Road Line"

    # Car icon at the LEFT of the slide
    car_body = add_car_shape(
        slide4,
        left=Inches(0.5),   # Left side of slide
        top=Inches(2.8),
        width=Inches(1.8),
        height=Inches(1.2)
    )

    # Start/End markers
    start_marker = slide4.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(1.0), Inches(0.4))
    start_marker.text_frame.paragraphs[0].text = "START"
    start_marker.text_frame.paragraphs[0].font.size = Pt(10)
    start_marker.text_frame.paragraphs[0].font.bold = True
    start_marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x80, 0x00)

    end_marker = slide4.shapes.add_textbox(Inches(8.5), Inches(4.2), Inches(1.0), Inches(0.4))
    end_marker.text_frame.paragraphs[0].text = "END"
    end_marker.text_frame.paragraphs[0].font.size = Pt(10)
    end_marker.text_frame.paragraphs[0].font.bold = True
    end_marker.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    # --- Slide 5: Conclusion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Expand EV fleet by 30% in Q2 2025"
    items5 = [
        "Implement predictive maintenance scheduling",
        "Deploy GPS-based route optimization across all depots",
        "Upgrade telematics systems in legacy vehicles",
        "Complete driver safety training certification program",
    ]
    for item in items5:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
