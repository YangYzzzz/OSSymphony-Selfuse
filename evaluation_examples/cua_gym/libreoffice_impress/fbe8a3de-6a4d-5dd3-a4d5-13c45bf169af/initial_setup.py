"""
Initial Setup: Create a 12-slide team meeting presentation with no transitions
Task ID: impress_tm_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_006'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=None, alignment=None):
    """Helper to add a formatted textbox."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q2 2025 Team Meeting"
    slide1.placeholders[1].text = "Engineering Division - April 15, 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "1. Project Status Updates"
    for item in ["2. Q1 Performance Review", "3. New Hiring Plan",
                 "4. Technical Debt Discussion", "5. Sprint Planning",
                 "6. Open Floor / Q&A"]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Project Alpha Status ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Project Alpha - Status"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Lead: Sarah Chen"
    for item in ["Milestone 3 completed on March 28",
                 "Performance benchmarks exceeded targets by 12%",
                 "Integration testing phase begins April 20",
                 "Risk: Dependency on external API (mitigation in progress)"]:
        p = body3.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 4: Project Beta Status ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Project Beta - Status"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Lead: Marcus Johnson"
    for item in ["User research complete - 47 interviews conducted",
                 "Design prototypes approved by stakeholders",
                 "Backend API development at 65% completion",
                 "On track for June 1 beta release"]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 5: Q1 Revenue Metrics ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Q1 Performance - Revenue"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Total Revenue: $2.4M (+18% YoY)"
    for item in ["SaaS Subscriptions: $1.8M (+22%)",
                 "Enterprise Deals: $420K (+8%)",
                 "Professional Services: $180K (-3%)",
                 "Customer Retention Rate: 94.2%"]:
        p = body5.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 6: Q1 Engineering Metrics ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Q1 Performance - Engineering"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Deployment Frequency: 4.2/week (up from 3.1)"
    for item in ["Mean Time to Recovery: 28 min (target: <30 min)",
                 "Code Coverage: 87% (+5% from Q4)",
                 "Bug Escape Rate: 2.1% (best quarter yet)",
                 "PRs Merged: 342 across 8 repositories"]:
        p = body6.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 7: Hiring Plan ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Q2 Hiring Plan"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Open Positions: 6"
    for item in ["2x Senior Backend Engineers (Project Alpha)",
                 "1x Staff Frontend Engineer (Project Beta)",
                 "1x DevOps Engineer (Platform Team)",
                 "1x QA Lead (Cross-team)",
                 "1x Product Designer (Growth Team)",
                 "Timeline: Offers out by May 30"]:
        p = body7.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 8: Technical Debt ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Technical Debt Priorities"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Critical Items (must address in Q2):"
    for item in ["Migrate authentication service to OAuth 2.1",
                 "Replace deprecated Redis 5.x cluster",
                 "Refactor payment processing module (>3 years old)",
                 "Update CI/CD pipeline to support ARM builds",
                 "Estimated effort: 180 engineering hours"]:
        p = body8.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 9: Sprint Planning ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Sprint 14 Planning"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Sprint Duration: April 14 - April 25"
    for item in ["Story Points Committed: 64",
                 "Carry-over from Sprint 13: 8 points",
                 "Focus: API Gateway refactor + Beta UI polish",
                 "Sprint Goal: Complete integration testing for Alpha M3"]:
        p = body9.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 10: Team Recognition ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Team Recognition"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Outstanding Contributions:"
    for item in ["Elena Rodriguez - Led critical production incident response",
                 "James Park - Mentored 3 junior developers this quarter",
                 "Priya Sharma - Delivered payment module 2 weeks early",
                 "David Kim - Improved build times by 40%"]:
        p = body10.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 11: Upcoming Events ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Upcoming Events"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "April - May 2025"
    for item in ["April 18: Company All-Hands (10:00 AM)",
                 "April 22-23: AWS Summit - 4 team members attending",
                 "May 2: Engineering Offsite (Lakeside Conference Center)",
                 "May 15: Q2 Mid-Quarter Review",
                 "May 20: Hackathon Week begins"]:
        p = body11.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 12: Q&A / Open Floor ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Questions & Open Discussion"
    body12 = slide12.placeholders[1].text_frame
    body12.text = "Thank you for your contributions this quarter!"
    p = body12.add_paragraph()
    p.text = ""
    p = body12.add_paragraph()
    p.text = "Next meeting: May 13, 2025 at 2:00 PM"
    p = body12.add_paragraph()
    p.text = "Action items will be shared via Slack #eng-team"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
