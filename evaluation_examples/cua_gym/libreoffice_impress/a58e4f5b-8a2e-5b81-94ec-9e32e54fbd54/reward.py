"""
FINAL REWARD SCRIPT - SUCCESS
Task: I've got a presentation where clicking 'Details' on the second slide should jump straight to a section marked 'Appendix A' on slide number 12. Can you guide me through setting up this kind of link in LibreOffice Impress?
Generated: 2025-08-07 09:36:54
Status: success
Model: o4-mini
Total Steps: 11
"""

#!/usr/bin/env python3
import os
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation

def verify_impress_task(file_path):
    print(f"Checking task completion for: {file_path}")
    score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: 0.0")
        return 0.0
    print(f"✓ File exists: (0.2 points)")
    score += 0.2

    # Requirement 2: Load presentation & slide count >=12 (0.2)
    try:
        prs = Presentation(file_path)
        total = len(prs.slides)
        print(f"Slide count: {total}")
        if total >= 12:
            print(f"✓ Slide count >=12 (0.2 points)")
            score += 0.2
        else:
            print(f"✗ Slide count <12")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print(f"REWARD: 0.0")
        return 0.0

    # Requirement 3: Find 'Details' shape on slide2 (0.2)
    details_shapes = []
    try:
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if hasattr(shape, 'text') and shape.text and 'details' in shape.text.lower().strip():
                details_shapes.append(shape)
                print(f"  ✓ Found shape with text '{shape.text.strip()}'")
        if details_shapes:
            print(f"✓ 'Details' shape found on slide2 (0.2 points)")
            score += 0.2
        else:
            print(f"✗ No 'Details' shape on slide2")
    except Exception as e:
        print(f"✗ Error checking slide2 shapes: {e}")

    # Requirement 4: Hyperlink exists on 'Details' shape (0.2)
    hlink_ids = []
    try:
        for shape in details_shapes:
            el = shape._element
            for h in el.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick'):
                rId = h.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                if rId:
                    hlink_ids.append(rId)
                    print(f"  ✓ Found hyperlink rId '{rId}' in 'Details' shape")
        if hlink_ids:
            print(f"✓ Hyperlink exists on 'Details' shape (0.2 points)")
            score += 0.2
        else:
            print(f"✗ No hyperlink found in 'Details' shape")
    except Exception as e:
        print(f"✗ Error checking hyperlink presence: {e}")

    # Requirement 5: Hyperlink points to slide12 (0.2)
    try:
        rels_path = 'ppt/slides/_rels/slide2.xml.rels'
        with zipfile.ZipFile(file_path, 'r') as z:
            data = z.read(rels_path)
            root = ET.fromstring(data)
            ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            found_target = False
            for rel in root.findall('r:Relationship', ns):
                rId = rel.get('Id')
                type_ = rel.get('Type')
                target = rel.get('Target')
                if rId in hlink_ids and type_.endswith('/hyperlink'):
                    print(f"    Relationship {rId}: target {target}, type {type_}")
                    if target.lower() == 'slide12':
                        print(f"    ✓ rId '{rId}' points to slide12")
                        found_target = True
                        break
            if found_target:
                print(f"✓ Hyperlink on 'Details' points to slide12 (0.2 points)")
                score += 0.2
            else:
                print(f"✗ Hyperlink on 'Details' does not point to slide12")
    except Exception as e:
        print(f"✗ Error checking hyperlink target: {e}")

    final = min(score, max_score)
    print(f"Total score: {final}/{max_score}")
    print(f"REWARD: {final}")
    return final

if __name__ == '__main__':
    path = '/home/user/ive_got_a_presentation_where_clicking_details_on_the_second_slide_should_jump_straight_to_a_section_.pptx'
    verify_impress_task(path)

