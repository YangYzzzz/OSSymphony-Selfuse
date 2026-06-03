"""
Initial Setup: Create Astronomy_101 presentation with 8 slides
Task ID: impress_teach_052
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_052'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, layout_idx, title_text, body_lines):
    """Helper to add a slide with title and bulleted body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Astronomy 101"
    slide1.placeholders[1].text = "A Journey Through the Cosmos"

    # Slide 2: Introduction
    add_title_body_slide(prs, 1, "Introduction to Astronomy", [
        "Astronomy is the study of celestial objects and phenomena",
        "Humans have gazed at the stars for over 5,000 years",
        "Modern astronomy uses telescopes, satellites, and spectroscopy",
        "This course covers our solar system and beyond",
    ])

    # Slide 3: The Solar System
    add_title_body_slide(prs, 1, "The Solar System", [
        "8 planets orbit our Sun in elliptical paths",
        "Inner planets: Mercury, Venus, Earth, Mars (rocky)",
        "Outer planets: Jupiter, Saturn, Uranus, Neptune (gas/ice giants)",
        "The asteroid belt separates the inner and outer planets",
        "Dwarf planets include Pluto, Eris, and Ceres",
    ])

    # Slide 4: Stars and Stellar Evolution
    add_title_body_slide(prs, 1, "Stars and Stellar Evolution", [
        "Stars form from collapsing clouds of gas and dust (nebulae)",
        "Nuclear fusion of hydrogen into helium powers main-sequence stars",
        "Our Sun is a G-type main-sequence star, about 4.6 billion years old",
        "Massive stars end their lives as supernovae, leaving neutron stars or black holes",
        "Low-mass stars become white dwarfs after shedding outer layers",
    ])

    # Slide 5: Galaxies
    add_title_body_slide(prs, 1, "Galaxies", [
        "Galaxies are massive systems of stars, gas, dust, and dark matter",
        "The Milky Way is a barred spiral galaxy ~100,000 light-years across",
        "Andromeda (M31) is our nearest large galaxy at 2.5 million light-years",
        "Galaxy types: spiral, elliptical, irregular, and lenticular",
        "The observable universe contains roughly 2 trillion galaxies",
    ])

    # Slide 6: Nebulae and Interstellar Medium
    add_title_body_slide(prs, 1, "Nebulae and the Interstellar Medium", [
        "Nebulae are vast clouds of gas and dust in interstellar space",
        "Emission nebulae glow from ionized hydrogen (e.g., Orion Nebula)",
        "Reflection nebulae scatter light from nearby stars",
        "Planetary nebulae form when dying stars expel their outer layers",
        "The interstellar medium fills the space between star systems",
    ])

    # Slide 7: Black Holes
    add_title_body_slide(prs, 1, "Black Holes", [
        "Black holes are regions where gravity prevents anything from escaping",
        "Stellar black holes form from collapsed massive stars (3-20 solar masses)",
        "Supermassive black holes reside at galaxy centers (millions to billions of solar masses)",
        "Sagittarius A* is the supermassive black hole at the Milky Way's center",
        "The Event Horizon Telescope captured the first image of a black hole in 2019",
    ])

    # Slide 8: Course Summary
    add_title_body_slide(prs, 1, "Course Summary", [
        "We explored the solar system, stars, galaxies, and cosmic phenomena",
        "Key concepts: fusion, gravity, light-years, spectral classification",
        "Astronomy continues to reveal new discoveries about our universe",
        "Next steps: observational techniques and telescope operation",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
