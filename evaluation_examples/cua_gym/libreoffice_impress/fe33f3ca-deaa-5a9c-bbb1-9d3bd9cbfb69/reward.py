"""
Reward Script: Quiz section insertion in Astronomy 101 presentation
Task ID: impress_teach_052
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Total slide count is 11 (3 new slides inserted)
  Component 2 (0.20): Slide 6 has 'Which planet is largest?' question with A-D options
  Component 3 (0.20): Slide 7 has 'What causes tides?' question with A-D options
  Component 4 (0.20): Slide 8 has 'The sun is a star: True or False?' with True/False options
  Component 5 (0.15): Each quiz slide has 'Click to Reveal' text box with white (hidden) text
  Component 6 (0.10): Original slides after quiz section are preserved (slides 9-11)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_052'


def get_all_text_from_slide(slide):
    """Get all text content from a slide, recursively handling groups."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def find_shape_with_text(slide, search_text):
    """Find a shape containing specific text (case-insensitive partial match)."""
    search_lower = search_text.lower()
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = " ".join(p.text for p in shape.text_frame.paragraphs).lower()
            if search_lower in full_text:
                return shape
    return None


def get_text_colors_in_shape(shape):
    """Get all font colors from runs in a shape. Returns list of color strings."""
    colors = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color.type is not None:
                        colors.append(str(run.font.color.rgb))
                    else:
                        colors.append("inherited")
                except Exception:
                    colors.append("unknown")
    return colors


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Total slide count is 11 (0.15 points)
    # Initial has 8 slides, golden should have 11 (3 new quiz slides inserted)
    try:
        if num_slides == 11:
            print(f"PASS: Component 1 — Slide count is 11 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 11 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # For remaining checks, we need at least 8 slides (quiz slides at positions 6-8)
    if num_slides < 8:
        print(f"CRITICAL: Not enough slides ({num_slides}) to check quiz content")
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 6 — 'Which planet is largest?' with A-D options (0.20 points)
    try:
        slide6 = prs.slides[5]  # 0-indexed
        all_text = get_all_text_from_slide(slide6)
        all_text_joined = " ".join(all_text).lower()

        has_question = "which planet is largest" in all_text_joined
        has_option_a = any("a)" in t.lower() for t in all_text)
        has_option_b = any("b)" in t.lower() for t in all_text)
        has_option_c = any("c)" in t.lower() for t in all_text)
        has_option_d = any("d)" in t.lower() for t in all_text)

        sub_score = 0.0
        if has_question:
            sub_score += 0.10
        if has_option_a and has_option_b and has_option_c and has_option_d:
            sub_score += 0.10

        if sub_score > 0:
            print(f"PASS: Component 2 — Slide 6 quiz question present (question={has_question}, options_ABCD={has_option_a and has_option_b and has_option_c and has_option_d}) ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 2 — Slide 6 missing quiz question or options. Text found: {all_text[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 7 — 'What causes tides?' with A-D options (0.20 points)
    try:
        slide7 = prs.slides[6]
        all_text = get_all_text_from_slide(slide7)
        all_text_joined = " ".join(all_text).lower()

        has_question = "what causes tides" in all_text_joined
        has_option_a = any("a)" in t.lower() for t in all_text)
        has_option_b = any("b)" in t.lower() for t in all_text)
        has_option_c = any("c)" in t.lower() for t in all_text)
        has_option_d = any("d)" in t.lower() for t in all_text)

        sub_score = 0.0
        if has_question:
            sub_score += 0.10
        if has_option_a and has_option_b and has_option_c and has_option_d:
            sub_score += 0.10

        if sub_score > 0:
            print(f"PASS: Component 3 — Slide 7 quiz question present (question={has_question}, options_ABCD={has_option_a and has_option_b and has_option_c and has_option_d}) ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 3 — Slide 7 missing quiz question or options. Text found: {all_text[:3]}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 8 — 'The sun is a star: True or False?' with True/False options (0.20 points)
    try:
        slide8 = prs.slides[7]
        all_text = get_all_text_from_slide(slide8)
        all_text_joined = " ".join(all_text).lower()

        has_question = "the sun is a star" in all_text_joined and "true or false" in all_text_joined
        has_true = any("true" in t.lower() for t in all_text)
        has_false = any("false" in t.lower() for t in all_text)

        sub_score = 0.0
        if has_question:
            sub_score += 0.10
        if has_true and has_false:
            sub_score += 0.10

        if sub_score > 0:
            print(f"PASS: Component 4 — Slide 8 T/F question present (question={has_question}, true={has_true}, false={has_false}) ({sub_score} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 4 — Slide 8 missing T/F question or options. Text found: {all_text[:3]}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Each quiz slide has 'Click to Reveal' box with white text (0.15 points)
    # White text = color FFFFFF, indicating hidden answer
    try:
        reveal_count = 0
        for si in [5, 6, 7]:  # slides 6, 7, 8 (0-indexed)
            slide = prs.slides[si]
            reveal_shape = find_shape_with_text(slide, "Click to Reveal")
            if reveal_shape is not None:
                colors = get_text_colors_in_shape(reveal_shape)
                # All text in reveal box should be white (FFFFFF)
                all_white = all(c == "FFFFFF" for c in colors if c not in ("inherited", "unknown"))
                if all_white and len([c for c in colors if c not in ("inherited", "unknown")]) > 0:
                    reveal_count += 1
                    print(f"  Slide {si+1}: 'Click to Reveal' found with white text")
                else:
                    print(f"  Slide {si+1}: 'Click to Reveal' found but text NOT white. Colors: {colors}")
            else:
                print(f"  Slide {si+1}: No 'Click to Reveal' shape found")

        if reveal_count == 3:
            print(f"PASS: Component 5 — All 3 quiz slides have hidden 'Click to Reveal' answer boxes (0.15 pts)")
            total_score += 0.15
        elif reveal_count > 0:
            partial = round(0.15 * reveal_count / 3, 2)
            print(f"PARTIAL: Component 5 — {reveal_count}/3 quiz slides have proper reveal boxes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No quiz slides have proper 'Click to Reveal' answer boxes")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Original slides preserved after quiz section (0.10 points)
    # Slides 9-11 should contain the same content as original slides 6-8
    # Original slide 6 title: "Nebulae and the Interstellar Medium"
    # Original slide 7 title: "Black Holes"
    # Original slide 8 title: "Course Summary"
    try:
        if num_slides >= 11:
            expected_titles = [
                (8, "Nebulae and the Interstellar Medium"),
                (9, "Black Holes"),
                (10, "Course Summary"),
            ]
            preserved_count = 0
            for idx, expected_title in expected_titles:
                slide = prs.slides[idx]
                slide_texts = get_all_text_from_slide(slide)
                slide_text_joined = " ".join(slide_texts).lower()
                if expected_title.lower() in slide_text_joined:
                    preserved_count += 1
                    print(f"  Slide {idx+1}: Found '{expected_title}'")
                else:
                    print(f"  Slide {idx+1}: Expected '{expected_title}', found texts: {slide_texts[:2]}")

            if preserved_count == 3:
                print(f"PASS: Component 6 — Original slides 6-8 preserved as slides 9-11 (0.10 pts)")
                total_score += 0.10
            elif preserved_count > 0:
                partial = round(0.10 * preserved_count / 3, 2)
                print(f"PARTIAL: Component 6 — {preserved_count}/3 original slides preserved ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 6 — Original slides not found at expected positions 9-11")
        else:
            print(f"FAIL: Component 6 — Not enough slides ({num_slides}) to check preservation")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
