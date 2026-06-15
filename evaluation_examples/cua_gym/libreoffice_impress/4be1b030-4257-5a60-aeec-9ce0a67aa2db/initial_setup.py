"""
Initial Setup: Create case study slide deck with 8 slides.
Slide 5 has 'Before & After' title with empty content.
Slide 6 has 'Client Testimonial' title with empty content.
Also creates placeholder images on Desktop and opens the file.
Task ID: impress_gf4_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
DESKTOP = '/home/user/Desktop'
TASK_ID = 'impress_gf4_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_placeholder_images():
    """Create before.png, after.png, and client_photo.png on the Desktop."""
    os.makedirs(DESKTOP, exist_ok=True)

    # before.png - a warehouse scene (simplified representation)
    img = Image.new('RGB', (800, 600), '#8B7355')
    draw = ImageDraw.Draw(img)
    # Draw shelves (messy/disorganized look)
    draw.rectangle([50, 100, 350, 550], fill='#A0522D', outline='#654321', width=2)
    draw.rectangle([400, 150, 750, 550], fill='#A0522D', outline='#654321', width=2)
    # Scattered boxes
    draw.rectangle([80, 200, 180, 280], fill='#DEB887', outline='#8B4513')
    draw.rectangle([200, 350, 320, 420], fill='#D2B48C', outline='#8B4513')
    draw.rectangle([420, 250, 550, 340], fill='#DEB887', outline='#8B4513')
    draw.rectangle([600, 400, 720, 500], fill='#D2B48C', outline='#8B4513')
    draw.rectangle([130, 430, 280, 530], fill='#F5DEB3', outline='#8B4513')
    # Title text
    draw.rectangle([150, 20, 650, 80], fill='#FFFFFF', outline='#333333')
    draw.text((220, 35), "Before: Warehouse Layout", fill='#333333')
    img.save(f'{DESKTOP}/before.png')

    # after.png - organized warehouse scene
    img = Image.new('RGB', (800, 600), '#E8E0D0')
    draw = ImageDraw.Draw(img)
    # Organized shelves
    for y in range(100, 500, 100):
        draw.rectangle([50, y, 350, y + 80], fill='#CD853F', outline='#8B4513', width=2)
        draw.rectangle([400, y, 750, y + 80], fill='#CD853F', outline='#8B4513', width=2)
    # Neatly placed boxes on shelves
    for y in range(100, 500, 100):
        for x in [70, 170, 270]:
            draw.rectangle([x, y + 10, x + 60, y + 60], fill='#4682B4', outline='#2F4F4F')
        for x in [420, 520, 640]:
            draw.rectangle([x, y + 10, x + 60, y + 60], fill='#4682B4', outline='#2F4F4F')
    draw.rectangle([150, 20, 650, 80], fill='#FFFFFF', outline='#333333')
    draw.text((220, 35), "After: Warehouse Layout", fill='#333333')
    img.save(f'{DESKTOP}/after.png')

    # client_photo.png - a stylized portrait placeholder
    img = Image.new('RGB', (400, 400), '#2C3E50')
    draw = ImageDraw.Draw(img)
    # Head
    draw.ellipse([140, 60, 260, 200], fill='#E8BEAC')
    # Body
    draw.ellipse([100, 200, 300, 380], fill='#3498DB')
    # Eyes
    draw.ellipse([170, 110, 190, 135], fill='#2C3E50')
    draw.ellipse([210, 110, 230, 135], fill='#2C3E50')
    # Smile
    draw.arc([175, 140, 225, 170], 0, 180, fill='#C0392B', width=2)
    img.save(f'{DESKTOP}/client_photo.png')

    print(f'Placeholder images created in {DESKTOP}/')


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Digital Transformation at Meridian Health"
    slide1.placeholders[1].text = "A Comprehensive Case Study\nPrepared by Apex Consulting Group\nMarch 2025"
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    for shape in slide1.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ---- Slide 2: Executive Summary ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.clear()
    summary_points = [
        "Meridian Health Systems operates 12 regional hospitals across the Pacific Northwest",
        "Legacy paper-based patient intake processes resulted in 45-minute average wait times",
        "Our digital transformation initiative reduced intake time by 68% within 6 months",
        "Annual cost savings of $2.3M from reduced administrative overhead",
        "Patient satisfaction scores improved from 3.2 to 4.7 out of 5.0",
    ]
    for i, point in enumerate(summary_points):
        if i == 0:
            body2.paragraphs[0].text = point
        else:
            p = body2.add_paragraph()
            p.text = point
            p.level = 0

    # ---- Slide 3: The Challenge ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "The Challenge"
    body3 = slide3.placeholders[1].text_frame
    body3.clear()
    challenges = [
        "Fragmented patient data across 5 incompatible legacy EHR systems",
        "Manual data entry leading to 12% error rate in patient records",
        "Staff spending 60% of time on administrative tasks vs. patient care",
        "No real-time visibility into bed availability or resource allocation",
        "Compliance risks with HIPAA due to inconsistent data handling",
        "Annual IT maintenance costs exceeding $4.8M for outdated systems",
    ]
    for i, ch in enumerate(challenges):
        if i == 0:
            body3.paragraphs[0].text = ch
        else:
            p = body3.add_paragraph()
            p.text = ch
            p.level = 0

    # ---- Slide 4: Solution Approach ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Solution Approach"
    body4 = slide4.placeholders[1].text_frame
    body4.clear()
    solutions = [
        "Phase 1: Unified cloud-based EHR platform deployment (Epic Systems)",
        "Phase 2: Automated patient intake kiosks with biometric verification",
        "Phase 3: Real-time dashboard for bed management and resource tracking",
        "Phase 4: AI-powered triage assistant for emergency department routing",
        "Integration with existing pharmacy and billing systems via HL7 FHIR APIs",
    ]
    for i, sol in enumerate(solutions):
        if i == 0:
            body4.paragraphs[0].text = sol
        else:
            p = body4.add_paragraph()
            p.text = sol
            p.level = 0

    # ---- Slide 5: Before & After (empty content area) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title text box
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf5 = title5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Before & After"
    p5.alignment = PP_ALIGN.CENTER
    run5 = p5.runs[0]
    run5.font.size = Pt(36)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    # Content area is intentionally empty — no images, shapes, or animations

    # ---- Slide 6: Client Testimonial (empty content area) ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf6 = title6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Client Testimonial"
    p6.alignment = PP_ALIGN.CENTER
    run6 = p6.runs[0]
    run6.font.size = Pt(36)
    run6.font.bold = True
    run6.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    # Content area is intentionally empty — no quotes, photos, or animations

    # ---- Slide 7: Results & Metrics ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Results & Metrics"
    body7 = slide7.placeholders[1].text_frame
    body7.clear()
    results = [
        "Patient intake time reduced from 45 minutes to 14 minutes (68% improvement)",
        "Data entry error rate dropped from 12% to 1.8%",
        "Staff time on administrative tasks reduced from 60% to 25%",
        "Real-time bed visibility achieved across all 12 facilities",
        "HIPAA compliance audit score improved from 72% to 98%",
        "ROI achieved within 18 months of initial deployment",
    ]
    for i, r in enumerate(results):
        if i == 0:
            body7.paragraphs[0].text = r
        else:
            p = body7.add_paragraph()
            p.text = r
            p.level = 0

    # ---- Slide 8: Next Steps & Contact ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps & Contact"
    body8 = slide8.placeholders[1].text_frame
    body8.clear()
    next_steps = [
        "Phase 5 planning: Telehealth integration across all campuses (Q3 2025)",
        "Expand AI triage to urgent care facilities",
        "Staff training program for advanced analytics dashboard",
        "Contact: Dr. Rachel Morrison, CTO — rachel.morrison@meridianhealth.org",
        "Project Lead: James Whitfield — james.whitfield@apexconsulting.com",
    ]
    for i, ns in enumerate(next_steps):
        if i == 0:
            body8.paragraphs[0].text = ns
        else:
            p = body8.add_paragraph()
            p.text = ns
            p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')


def main():
    create_placeholder_images()
    create_presentation()
    # Open in LibreOffice Impress for GUI-ready state
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
