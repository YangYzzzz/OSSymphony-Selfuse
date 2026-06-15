"""
Reward Script: Visually rich case study slide deck with before/after comparison and animated testimonial
Task ID: impress_gf4_037
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 5 has two images (before.png and after.png) side by side (0.20)
  Component 2: Slide 5 has descriptive text labels for the images (0.10)
  Component 3: Slide 5 has a Split transition (0.15)
  Component 4: Slide 6 has a large quote text box (0.15)
  Component 5: Slide 6 has an oval/circle shape with picture fill (client photo) (0.15)
  Component 6: Slide 6 has speaker name and title text (0.10)
  Component 7: Slide 6 has animations on at least 3 shapes (0.15)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_037'


def persist_app_state(domain):
    """Try to save any open LibreOffice document."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]
    slide6 = prs.slides[5]

    # =========================================================================
    # Component 1: Slide 5 has two image objects (before and after) side by side (0.20 points)
    # =========================================================================
    try:
        pictures_s5 = [s for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if len(pictures_s5) >= 2:
            # Check they are side-by-side: one should be on the left half, one on the right half
            # Or at minimum, their left positions should differ significantly
            lefts = sorted([p.left for p in pictures_s5])
            # Side by side means they don't fully overlap horizontally
            # The rightmost left position should be > leftmost left + leftmost width * 0.5
            pic_by_left = sorted(pictures_s5, key=lambda s: s.left)
            left_pic = pic_by_left[0]
            right_pic = pic_by_left[1]
            horizontal_gap = right_pic.left - (left_pic.left + left_pic.width)
            # They should be roughly side by side: right pic starts near or after left pic ends
            if right_pic.left > left_pic.left + left_pic.width * 0.3:
                print(f"PASS: Component 1 — Slide 5 has {len(pictures_s5)} images side by side (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 — Images not side by side. Left positions: {left_pic.left}, {right_pic.left}")
        else:
            print(f"FAIL: Component 1 — Slide 5 has {len(pictures_s5)} pictures, need at least 2")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # =========================================================================
    # Component 2: Slide 5 has descriptive text labels (not just the title) (0.10 points)
    # =========================================================================
    try:
        # Count text shapes on slide 5 that are not the title and not just "Before & After"
        text_shapes_s5 = []
        for shape in slide5.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text.lower() != 'before & after':
                    text_shapes_s5.append(text)

        # Need at least 2 descriptive text labels (one for before, one for after)
        descriptive_labels = [t for t in text_shapes_s5 if len(t) > 5]
        if len(descriptive_labels) >= 2:
            print(f"PASS: Component 2 — Slide 5 has {len(descriptive_labels)} descriptive text labels (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 5 has {len(descriptive_labels)} descriptive text labels, need >= 2")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =========================================================================
    # Component 3: Slide 5 has a Split transition (0.15 points)
    # =========================================================================
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                root = ET.parse(f).getroot()
                tr = root.find('.//p:transition', ns)
                if tr is not None:
                    split_el = tr.find('.//p:split', ns)
                    if split_el is not None:
                        print(f"PASS: Component 3 — Slide 5 has Split transition (0.15 pts)")
                        total_score += 0.15
                    else:
                        # Check all child elements for split-like transitions
                        children = list(tr)
                        child_tags = [c.tag.split('}')[-1] if '}' in c.tag else c.tag for c in children]
                        print(f"FAIL: Component 3 — Slide 5 has transition but not Split. Found: {child_tags}")
                else:
                    print(f"FAIL: Component 3 — Slide 5 has no transition")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =========================================================================
    # Component 4: Slide 6 has a large quote text box (0.15 points)
    # =========================================================================
    try:
        quote_found = False
        for shape in slide6.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Quote should be a substantial text, likely starting with quote mark or
                # containing testimonial-like content. Must be longer than simple labels.
                if len(text) > 50 and text != 'Client Testimonial':
                    quote_found = True
                    print(f"PASS: Component 4 — Slide 6 has quote text: '{text[:60]}...' (0.15 pts)")
                    total_score += 0.15
                    break
        if not quote_found:
            print(f"FAIL: Component 4 — No large quote text found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =========================================================================
    # Component 5: Slide 6 has an oval/circle shape with picture fill (client photo) (0.15 points)
    # =========================================================================
    try:
        oval_with_pic = False
        for shape in slide6.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                try:
                    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                    if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                        # Check for picture fill
                        fill = shape.fill
                        if fill.type is not None and fill.type == 6:  # PICTURE fill type
                            oval_with_pic = True
                            print(f"PASS: Component 5 — Slide 6 has oval shape with picture fill (0.15 pts)")
                            total_score += 0.15
                            break
                except Exception:
                    pass

            # Also check for a picture shape that might be oval-cropped
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # If there's a picture on slide 6, it could be the client photo
                # Check if it has any oval crop or frame
                oval_with_pic = True
                print(f"PASS: Component 5 — Slide 6 has picture shape (client photo) (0.15 pts)")
                total_score += 0.15
                break

        if not oval_with_pic:
            # Also check via XML for blipFill inside oval shapes
            try:
                ns_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slides/slide6.xml') as f:
                        root = ET.parse(f).getroot()
                        # Find any shape with oval preset and blipFill
                        for sp in root.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sp', {}):
                            xml_str = ET.tostring(sp, encoding='unicode')
                            if 'ellipse' in xml_str.lower() or 'oval' in xml_str.lower():
                                if 'blipFill' in xml_str or 'blip' in xml_str:
                                    oval_with_pic = True
                                    print(f"PASS: Component 5 — Slide 6 has oval with picture via XML (0.15 pts)")
                                    total_score += 0.15
                                    break
            except Exception:
                pass

        if not oval_with_pic:
            print(f"FAIL: Component 5 — No oval/circle shape with picture fill found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # =========================================================================
    # Component 6: Slide 6 has speaker name and title text (0.10 points)
    # =========================================================================
    try:
        speaker_info_found = False
        for shape in slide6.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip().lower()
                # Look for a shape containing name-like content (not the quote, not the header)
                # Should contain something like a person name and a title/position
                if 'client testimonial' in full_text:
                    continue
                if len(full_text) > 200:
                    continue  # Skip the quote text
                # Check for typical speaker info patterns
                has_name = any(word in full_text for word in ['dr.', 'mr.', 'ms.', 'mrs.', 'morrison', 'rachel'])
                has_title = any(word in full_text for word in ['officer', 'director', 'manager', 'cto', 'ceo', 'chief', 'president', 'meridian'])
                if has_name or has_title:
                    speaker_info_found = True
                    print(f"PASS: Component 6 — Slide 6 has speaker info: '{shape.text_frame.text.strip()[:60]}' (0.10 pts)")
                    total_score += 0.10
                    break
        if not speaker_info_found:
            print(f"FAIL: Component 6 — No speaker name/title text found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # =========================================================================
    # Component 7: Slide 6 has animations on at least 3 shapes (0.15 points)
    # =========================================================================
    try:
        ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide6.xml') as f:
                root = ET.parse(f).getroot()
                timing = root.find('.//p:timing', ns)
                if timing is not None:
                    # Count unique animated shape targets
                    animated_spids = set()
                    for spTgt in timing.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spTgt'):
                        spid = spTgt.get('spid')
                        if spid:
                            animated_spids.add(spid)

                    # Also check for animation effects (fade, motion, etc.)
                    anim_effects = timing.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}animEffect')
                    anim_motions = timing.findall('.//{http://schemas.openxmlformats.org/presentationml/2006/main}animMotion')
                    total_anims = len(anim_effects) + len(anim_motions)

                    if len(animated_spids) >= 3:
                        print(f"PASS: Component 7 — Slide 6 has animations on {len(animated_spids)} shapes ({total_anims} effects) (0.15 pts)")
                        total_score += 0.15
                    elif len(animated_spids) >= 2:
                        print(f"PARTIAL: Component 7 — Slide 6 has animations on {len(animated_spids)} shapes, need 3 (0.07 pts)")
                        total_score += 0.07
                    else:
                        print(f"FAIL: Component 7 — Slide 6 has animations on {len(animated_spids)} shapes, need at least 3")
                else:
                    print(f"FAIL: Component 7 — No animations found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
