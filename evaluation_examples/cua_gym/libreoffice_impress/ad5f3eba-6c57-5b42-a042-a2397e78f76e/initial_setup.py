"""
Initial Setup: Create a 30-slide Manual Training presentation without page numbering.
Task ID: impress_ma_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, alignment=PP_ALIGN.LEFT, color=None):
    """Helper to add a formatted text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14):
    """Helper to add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(font_size)
    return txBox


def create_initial():
    prs = Presentation()

    # Standard widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content for a 30-slide Manual Training presentation
    slide_data = [
        # Slide 1: Title
        {"layout": 0, "title": "Employee Onboarding Manual",
         "subtitle": "Comprehensive Training Program 2025\nHR Department — Nextera Solutions Inc."},
        # Slide 2: Agenda
        {"layout": 1, "title": "Training Agenda",
         "bullets": ["Company Overview & Culture", "Workplace Policies & Compliance",
                      "IT Systems & Security Protocols", "Role-Specific Training Modules",
                      "Performance Review Process", "Benefits & Compensation",
                      "Health & Safety Guidelines", "Q&A and Wrap-Up"]},
        # Slide 3
        {"layout": 1, "title": "About Nextera Solutions",
         "bullets": ["Founded in 2008 in San Francisco, CA",
                      "Over 2,400 employees across 12 global offices",
                      "Revenue: $890M (FY2024)", "Industry leader in enterprise SaaS platforms",
                      "Core values: Innovation, Integrity, Collaboration"]},
        # Slide 4
        {"layout": 1, "title": "Our Leadership Team",
         "bullets": ["CEO: Dr. Amanda Richardson", "CTO: Rajesh Patel",
                      "CFO: Michael Torres", "VP Engineering: Sarah Kim",
                      "VP Product: David Okonkwo", "Chief People Officer: Lisa Nakamura"]},
        # Slide 5
        {"layout": 1, "title": "Company Mission & Vision",
         "bullets": ["Mission: Empower businesses with intelligent automation",
                      "Vision: A world where technology amplifies human potential",
                      "Strategy 2025-2028: Expand AI-driven product suite",
                      "Focus areas: Healthcare, Finance, Education"]},
        # Slide 6
        {"layout": 1, "title": "Workplace Code of Conduct",
         "bullets": ["Respect and professionalism at all times",
                      "Zero tolerance for harassment or discrimination",
                      "Confidential information must be protected",
                      "Report concerns via Ethics Hotline: ext. 4500",
                      "Annual compliance training required by March 31"]},
        # Slide 7
        {"layout": 1, "title": "Communication Guidelines",
         "bullets": ["Primary channels: Slack, Email, Zoom",
                      "Response time expectations: 4 hours internal, 2 hours client",
                      "All-hands meetings: First Monday of each month",
                      "Team standups: Daily at 9:30 AM local time",
                      "Document decisions in Confluence"]},
        # Slide 8
        {"layout": 1, "title": "IT Systems Overview",
         "bullets": ["Laptop setup: MacBook Pro or Dell XPS (your choice)",
                      "VPN required for remote access: GlobalProtect",
                      "Password policy: 12+ chars, rotated every 90 days",
                      "Two-factor authentication mandatory",
                      "Software requests via ServiceNow portal"]},
        # Slide 9
        {"layout": 1, "title": "Data Security Protocols",
         "bullets": ["Classify data: Public, Internal, Confidential, Restricted",
                      "Never share credentials via email or chat",
                      "Encrypt sensitive files before transmission",
                      "Report phishing attempts to security@nextera.com",
                      "USB drives prohibited on company devices"]},
        # Slide 10
        {"layout": 1, "title": "Access & Permissions",
         "bullets": ["LDAP-based role access control",
                      "Request access via IT ticket within first 48 hours",
                      "Manager approval required for elevated permissions",
                      "Quarterly access reviews by department leads",
                      "Offboarding revokes all access within 24 hours"]},
        # Slide 11
        {"layout": 1, "title": "Engineering Onboarding",
         "bullets": ["Clone repos from GitHub Enterprise",
                      "Dev environment setup guide: wiki/eng-setup",
                      "Code review required: minimum 2 approvals",
                      "CI/CD pipeline: Jenkins + ArgoCD",
                      "Sprint cadence: 2-week cycles, planning on Mondays"]},
        # Slide 12
        {"layout": 1, "title": "Product Development Lifecycle",
         "bullets": ["Phase 1: Discovery & Research (2-3 weeks)",
                      "Phase 2: Design & Prototyping (2 weeks)",
                      "Phase 3: Development & Testing (4-6 weeks)",
                      "Phase 4: Staging & QA (1-2 weeks)",
                      "Phase 5: Production Release & Monitoring"]},
        # Slide 13
        {"layout": 1, "title": "Quality Assurance Standards",
         "bullets": ["Unit test coverage minimum: 80%",
                      "Integration tests run nightly",
                      "Performance benchmarks tracked in Grafana",
                      "Bug severity levels: P0 (critical) to P3 (minor)",
                      "P0 bugs require fix within 4 hours"]},
        # Slide 14
        {"layout": 1, "title": "Sales & Marketing Training",
         "bullets": ["CRM: Salesforce — log all client interactions",
                      "Lead scoring model: BANT framework",
                      "Quarterly pipeline reviews with VP Sales",
                      "Marketing campaigns coordinated via HubSpot",
                      "Brand guidelines available at brand.nextera.com"]},
        # Slide 15
        {"layout": 1, "title": "Client Engagement Protocol",
         "bullets": ["Initial outreach within 24 hours of lead assignment",
                      "Discovery call template in Notion",
                      "Proposal turnaround: 5 business days maximum",
                      "Contract review: Legal team SLA is 48 hours",
                      "NPS survey sent 30 days post-deployment"]},
        # Slide 16
        {"layout": 1, "title": "Performance Review Framework",
         "bullets": ["Review cycles: Semi-annual (June & December)",
                      "Self-assessment due 2 weeks before review",
                      "Manager calibration sessions in week 3",
                      "Rating scale: Exceeds / Meets / Below expectations",
                      "Development plan required for all ratings"]},
        # Slide 17
        {"layout": 1, "title": "Goal Setting with OKRs",
         "bullets": ["Company OKRs cascade to team and individual level",
                      "3-5 objectives per quarter recommended",
                      "Key results must be measurable and time-bound",
                      "Mid-quarter check-in with manager required",
                      "OKR tracking in Lattice platform"]},
        # Slide 18
        {"layout": 1, "title": "Career Development Paths",
         "bullets": ["Individual Contributor track: IC1 through IC6",
                      "Management track: M1 (Team Lead) through M5 (VP)",
                      "Lateral moves encouraged after 18 months",
                      "Mentorship program: sign up in People portal",
                      "Learning budget: $2,500/year per employee"]},
        # Slide 19
        {"layout": 1, "title": "Compensation & Benefits Overview",
         "bullets": ["Salary bands reviewed annually against market data",
                      "Annual bonus: 10-20% based on performance",
                      "Equity grants vest over 4 years with 1-year cliff",
                      "401(k) match: 100% up to 6% of salary",
                      "Employee stock purchase plan at 15% discount"]},
        # Slide 20
        {"layout": 1, "title": "Health & Wellness Benefits",
         "bullets": ["Medical: Blue Cross PPO and HMO options",
                      "Dental & Vision: Delta Dental, VSP",
                      "Mental health: 12 free therapy sessions/year via Lyra",
                      "Gym reimbursement: up to $100/month",
                      "Wellness days: 2 additional PTO days per quarter"]},
        # Slide 21
        {"layout": 1, "title": "Time Off & Leave Policies",
         "bullets": ["PTO: 20 days/year for new employees, increases to 25 at year 3",
                      "Sick leave: 10 days/year (no carryover)",
                      "Parental leave: 16 weeks paid for all parents",
                      "Bereavement: 5 days for immediate family",
                      "Jury duty: full pay for up to 10 days"]},
        # Slide 22
        {"layout": 1, "title": "Remote Work Guidelines",
         "bullets": ["Hybrid model: 3 days in-office, 2 days remote",
                      "Core collaboration hours: 10 AM - 3 PM local",
                      "Home office stipend: $500 one-time setup",
                      "Reliable internet (50 Mbps+) required",
                      "International remote work requires VP approval"]},
        # Slide 23
        {"layout": 1, "title": "Health & Safety in the Office",
         "bullets": ["Emergency exits marked on every floor plan",
                      "Fire drills conducted quarterly",
                      "First aid kits in each kitchen area",
                      "Ergonomic assessment available via Facilities",
                      "Report hazards to safety@nextera.com immediately"]},
        # Slide 24
        {"layout": 1, "title": "Incident Response Procedures",
         "bullets": ["Step 1: Ensure personal safety",
                      "Step 2: Contact Security (ext. 9111)",
                      "Step 3: Document incident in SafetyNet portal",
                      "Step 4: Manager notification within 1 hour",
                      "Step 5: Follow-up investigation within 48 hours"]},
        # Slide 25
        {"layout": 1, "title": "Diversity, Equity & Inclusion",
         "bullets": ["Employee Resource Groups: 8 active groups",
                      "Unconscious bias training: mandatory annually",
                      "Pay equity audits conducted bi-annually",
                      "Diverse interview panels required for all roles",
                      "DEI metrics reported to Board quarterly"]},
        # Slide 26
        {"layout": 1, "title": "Corporate Social Responsibility",
         "bullets": ["Volunteer time: 16 hours/year paid",
                      "Matching donations: up to $1,000/year",
                      "Carbon neutral commitment by 2027",
                      "Annual sustainability report published in Q1",
                      "Community partnerships in each office city"]},
        # Slide 27
        {"layout": 1, "title": "Expense & Travel Policy",
         "bullets": ["Expense reports via Concur within 30 days",
                      "Flights: economy class for trips under 6 hours",
                      "Hotel limit: $250/night in major metros",
                      "Meals: $75/day domestic, $100/day international",
                      "Manager pre-approval for trips over $2,000"]},
        # Slide 28
        {"layout": 1, "title": "Internal Tools & Resources",
         "bullets": ["Confluence: documentation & knowledge base",
                      "Jira: project & task management",
                      "Notion: team wikis and templates",
                      "Figma: design collaboration",
                      "Datadog: monitoring and observability"]},
        # Slide 29
        {"layout": 1, "title": "Key Contacts & Support",
         "bullets": ["IT Help Desk: helpdesk@nextera.com / ext. 5000",
                      "HR Business Partner: hr-support@nextera.com",
                      "Facilities: facilities@nextera.com / ext. 3200",
                      "Ethics Hotline: anonymous, ext. 4500",
                      "Your Manager: scheduled 1:1 every week"]},
        # Slide 30
        {"layout": 1, "title": "Thank You & Next Steps",
         "bullets": ["Complete onboarding checklist by end of Week 1",
                      "Schedule meet-and-greets with key stakeholders",
                      "Set up 30/60/90 day goals with your manager",
                      "Join your team Slack channel",
                      "Welcome to Nextera Solutions — we're glad you're here!"]},
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

        # Set subtitle or bullets
        if "subtitle" in sd:
            if 1 in slide.placeholders:
                slide.placeholders[1].text = sd["subtitle"]
        elif "bullets" in sd:
            if 1 in slide.placeholders:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for j, bullet in enumerate(sd["bullets"]):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(16)
                        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
