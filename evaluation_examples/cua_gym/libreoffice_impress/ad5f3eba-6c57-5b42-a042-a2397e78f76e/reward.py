"""
Reward Script: Add 'Page X of Y' text box to bottom center of master slide
Task ID: impress_ma_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Text box on master with slidenum field
  Component 2 (0.25): Text contains 'Page' prefix before slide number
  Component 3 (0.20): Text box is horizontally centered on slide
  Component 4 (0.20): Text box is at bottom of slide with small font
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_028'


def persist_app_state(domain: str):
    """Save any unsaved LibreOffice edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the master slide has a new text box at bottom center
    containing 'Page <slide#>' with a slidenum field.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Cannot import python-pptx: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get the first slide master
    try:
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find text boxes on the master that contain a slidenum field
    # We need to distinguish from the default "Slide Number Placeholder 5"
    # which is a PLACEHOLDER type, not a TEXT_BOX
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

    candidate_shapes = []
    for shape in master.shapes:
        if not shape.has_text_frame:
            continue
        xml_str = shape.text_frame._txBody.xml
        # Look for slidenum field in any non-placeholder shape OR in any shape with "Page" text
        has_slidenum = 'type="slidenum"' in xml_str
        has_page_text = 'page' in shape.text_frame.text.lower()
        is_textbox = (shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX)
        # Accept text boxes with slidenum, or any shape that has both "Page" and slidenum
        if has_slidenum and (is_textbox or has_page_text):
            # Exclude the default slide number placeholder (which just has <#> without "Page")
            if is_textbox or has_page_text:
                candidate_shapes.append(shape)

    # Component 1: A text box on the master contains a slidenum field (0.35 points)
    try:
        if len(candidate_shapes) > 0:
            # Found at least one matching shape
            target = candidate_shapes[0]
            xml_str = target.text_frame._txBody.xml
            if 'type="slidenum"' in xml_str:
                print(f"PASS: Component 1 — Found text box '{target.name}' on master with slidenum field (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 — Shape '{target.name}' found but no slidenum field")
        else:
            # Also check via XML in the slideMaster zip for any non-placeholder text box with slidenum
            found_via_xml = False
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
                        root = ET.parse(f).getroot()
                        ns = {
                            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        }
                        # Look for sp elements that have a:fld with type="slidenum"
                        for sp in root.findall('.//p:sp', ns):
                            flds = sp.findall('.//a:fld', ns)
                            for fld in flds:
                                if fld.get('type') == 'slidenum':
                                    # Check if it's NOT a placeholder (no nvSpPr/nvPr/ph)
                                    nvPr = sp.find('.//p:nvSpPr/p:nvPr', ns)
                                    ph = nvPr.find('p:ph', ns) if nvPr is not None else None
                                    # Also check text content for "Page"
                                    texts = [t.text or '' for t in sp.findall('.//a:t', ns)]
                                    full_text = ''.join(texts).lower()
                                    if ph is None or 'page' in full_text:
                                        found_via_xml = True
                                        break
                            if found_via_xml:
                                break
            except Exception as e:
                print(f"  XML fallback error: {e}")

            if found_via_xml:
                print(f"PASS: Component 1 — Found text box on master with slidenum field via XML (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 1 — No text box with slidenum field found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: The text contains "Page" prefix (0.25 points)
    try:
        found_page_prefix = False
        # Check via python-pptx candidate shapes
        for shape in candidate_shapes:
            text = shape.text_frame.text.lower()
            if 'page' in text:
                found_page_prefix = True
                print(f"PASS: Component 2 — Text contains 'Page' prefix: '{shape.text_frame.text}' (0.25 pts)")
                total_score += 0.25
                break

        if not found_page_prefix:
            # Check via ZIP XML
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
                        root = ET.parse(f).getroot()
                        ns = {
                            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        }
                        for sp in root.findall('.//p:sp', ns):
                            flds = sp.findall('.//a:fld', ns)
                            has_slidenum = any(f.get('type') == 'slidenum' for f in flds)
                            if has_slidenum:
                                texts = [t.text or '' for t in sp.findall('.//a:t', ns)]
                                full_text = ''.join(texts).lower()
                                if 'page' in full_text:
                                    found_page_prefix = True
                                    print(f"PASS: Component 2 — Text contains 'Page' prefix via XML (0.25 pts)")
                                    total_score += 0.25
                                    break
            except Exception as e:
                print(f"  XML fallback error: {e}")

            if not found_page_prefix:
                print(f"FAIL: Component 2 — No 'Page' prefix found in slide number text box")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Text box is horizontally centered on the slide (0.20 points)
    try:
        centered = False
        slide_width = prs.slide_width

        for shape in candidate_shapes:
            shape_center_x = shape.left + shape.width // 2
            slide_center_x = slide_width // 2
            # Allow 10% tolerance for centering
            tolerance = slide_width * 0.10
            offset = abs(shape_center_x - slide_center_x)
            if offset <= tolerance:
                centered = True
                print(f"PASS: Component 3 — Text box is centered (offset: {offset} EMU, tolerance: {tolerance} EMU) (0.20 pts)")
                total_score += 0.20
                break

        if not centered and len(candidate_shapes) == 0:
            # Try XML approach to get position
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
                        root = ET.parse(f).getroot()
                        ns = {
                            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        }
                        for sp in root.findall('.//p:sp', ns):
                            flds = sp.findall('.//a:fld', ns)
                            has_slidenum = any(f.get('type') == 'slidenum' for f in flds)
                            texts = [t.text or '' for t in sp.findall('.//a:t', ns)]
                            full_text = ''.join(texts).lower()
                            if has_slidenum and 'page' in full_text:
                                off_elem = sp.find('.//p:spPr/a:xfrm/a:off', ns)
                                ext_elem = sp.find('.//p:spPr/a:xfrm/a:ext', ns)
                                if off_elem is not None and ext_elem is not None:
                                    x = int(off_elem.get('x', 0))
                                    cx = int(ext_elem.get('cx', 0))
                                    shape_center = x + cx // 2
                                    slide_cx = slide_width // 2
                                    tol = slide_width * 0.10
                                    if abs(shape_center - slide_cx) <= tol:
                                        centered = True
                                        print(f"PASS: Component 3 — Text box centered via XML (0.20 pts)")
                                        total_score += 0.20
                                break
            except Exception as e:
                print(f"  XML fallback error: {e}")

        if not centered:
            print(f"FAIL: Component 3 — Text box is not centered horizontally")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Text box is at the bottom of the slide with small font (0.20 points)
    try:
        at_bottom = False
        slide_height = prs.slide_height

        for shape in candidate_shapes:
            # "Bottom" means the shape's top is in the lower 25% of the slide
            bottom_threshold = slide_height * 0.75
            if shape.top >= bottom_threshold:
                # Check font size is small (<=14pt = 177800 EMU; or 1000 centipoints in sz attr)
                font_ok = False
                xml_str = shape.text_frame._txBody.xml
                # Parse sz attribute from rPr
                import re
                sz_matches = re.findall(r'sz="(\d+)"', xml_str)
                if sz_matches:
                    # sz is in hundredths of a point
                    max_sz = max(int(s) for s in sz_matches)
                    if max_sz <= 1400:  # 14pt or less
                        font_ok = True

                if font_ok:
                    at_bottom = True
                    print(f"PASS: Component 4 — Text box at bottom (top={shape.top}, threshold={bottom_threshold}) with small font (0.20 pts)")
                    total_score += 0.20
                    break
                else:
                    # Still give partial if position is correct but font is large
                    at_bottom = True
                    print(f"PARTIAL: Component 4 — Text box at bottom but font may be large (0.10 pts)")
                    total_score += 0.10
                    break

        if not at_bottom and len(candidate_shapes) == 0:
            # XML fallback
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
                        root = ET.parse(f).getroot()
                        ns = {
                            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                        }
                        for sp in root.findall('.//p:sp', ns):
                            flds = sp.findall('.//a:fld', ns)
                            has_slidenum = any(f.get('type') == 'slidenum' for f in flds)
                            texts = [t.text or '' for t in sp.findall('.//a:t', ns)]
                            full_text = ''.join(texts).lower()
                            if has_slidenum and 'page' in full_text:
                                off_elem = sp.find('.//p:spPr/a:xfrm/a:off', ns)
                                if off_elem is not None:
                                    y = int(off_elem.get('y', 0))
                                    if y >= slide_height * 0.75:
                                        at_bottom = True
                                        print(f"PASS: Component 4 — Text box at bottom via XML (0.20 pts)")
                                        total_score += 0.20
                                break
            except Exception as e:
                print(f"  XML fallback error: {e}")

        if not at_bottom:
            print(f"FAIL: Component 4 — Text box not at bottom of slide")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
