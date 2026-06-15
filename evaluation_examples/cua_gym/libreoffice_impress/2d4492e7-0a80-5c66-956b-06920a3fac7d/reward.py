"""
Reward Script: Make all text in the presentation bold.
Task ID: impstruct_016
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 1 — all text runs are bold (0.35 pts)
  Component 2: Slide 2 — all text runs are bold (0.35 pts)
  Component 3: Slide 3 — all text runs are bold (0.30 pts)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impstruct_016'


def get_all_nonempty_runs(slide):
    """Collect all non-empty text runs from all text-bearing shapes on a slide,
    including nested shapes inside groups."""
    runs = []

    def extract_runs(shape):
        if hasattr(shape, 'text_frame'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or '').strip():
                        runs.append(run)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                extract_runs(sub)

    for shape in slide.shapes:
        extract_runs(shape)
    return runs


def check_all_bold(runs):
    """Return True if every run has bold == True.
    Per python-pptx, bold can be True, False, or None (inherit).
    Only explicit True counts as bold."""
    if not runs:
        return False  # no text to verify — cannot confirm bold
    for run in runs:
        if run.font.bold is not True:
            return False
    return True


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    if len(slides) < 3:
        print(f"FAIL: Expected at least 3 slides, found {len(slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 1 — all text runs are bold (0.35 points)
    try:
        runs_s1 = get_all_nonempty_runs(slides[0])
        if check_all_bold(runs_s1):
            print(f"PASS: Component 1 — Slide 1: all {len(runs_s1)} runs are bold (0.35 pts)")
            total_score += 0.35
        else:
            non_bold = [r.text for r in runs_s1 if r.font.bold is not True]
            print(f"FAIL: Component 1 — Slide 1: {len(non_bold)} of {len(runs_s1)} runs are not bold: {non_bold[:3]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 — all text runs are bold (0.35 points)
    try:
        runs_s2 = get_all_nonempty_runs(slides[1])
        if check_all_bold(runs_s2):
            print(f"PASS: Component 2 — Slide 2: all {len(runs_s2)} runs are bold (0.35 pts)")
            total_score += 0.35
        else:
            non_bold = [r.text for r in runs_s2 if r.font.bold is not True]
            print(f"FAIL: Component 2 — Slide 2: {len(non_bold)} of {len(runs_s2)} runs are not bold: {non_bold[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 — all text runs are bold (0.30 points)
    try:
        runs_s3 = get_all_nonempty_runs(slides[2])
        if check_all_bold(runs_s3):
            print(f"PASS: Component 3 — Slide 3: all {len(runs_s3)} runs are bold (0.30 pts)")
            total_score += 0.30
        else:
            non_bold = [r.text for r in runs_s3 if r.font.bold is not True]
            print(f"FAIL: Component 3 — Slide 3: {len(non_bold)} of {len(runs_s3)} runs are not bold: {non_bold[:3]}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
