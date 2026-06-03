"""
Initial Setup: Make all text in the presentation bold.
Task ID: impstruct_016
Domain: libreoffice_impress

Creates a 3-slide presentation with titles and body text, none bold.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_016'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_textframe(tf, text, font_name="Calibri", font_size=Pt(18)):
    """Add text to a text frame with explicit non-bold formatting."""
    p = tf.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = False
    return run


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    # Title
    title_tf = slide1.shapes.title.text_frame
    title_tf.paragraphs[0].text = ""
    title_run = title_tf.paragraphs[0].add_run()
    title_run.text = "Q3 Strategy Review"
    title_run.font.name = "Calibri"
    title_run.font.size = Pt(36)
    title_run.font.bold = False

    # Subtitle
    subtitle_tf = slide1.placeholders[1].text_frame
    subtitle_tf.paragraphs[0].text = ""
    sub_run = subtitle_tf.paragraphs[0].add_run()
    sub_run.text = "Key insights from the quarterly planning session"
    sub_run.font.name = "Calibri"
    sub_run.font.size = Pt(20)
    sub_run.font.bold = False

    # --- Slide 2: Title + Content ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    # Title
    title_tf2 = slide2.shapes.title.text_frame
    title_tf2.paragraphs[0].text = ""
    t2_run = title_tf2.paragraphs[0].add_run()
    t2_run.text = "Revenue Highlights"
    t2_run.font.name = "Calibri"
    t2_run.font.size = Pt(32)
    t2_run.font.bold = False

    # Body content
    body_tf2 = slide2.placeholders[1].text_frame
    body_tf2.paragraphs[0].text = ""

    body_items = [
        "Total revenue reached $4.2M, up 18% from Q2",
        "Enterprise segment grew 25% quarter-over-quarter",
        "New customer acquisition cost decreased by $340",
        "Subscription renewals maintained 94% retention rate",
    ]
    for i, item in enumerate(body_items):
        if i == 0:
            p = body_tf2.paragraphs[0]
        else:
            p = body_tf2.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.bold = False

    # --- Slide 3: Title + Content ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    # Title
    title_tf3 = slide3.shapes.title.text_frame
    title_tf3.paragraphs[0].text = ""
    t3_run = title_tf3.paragraphs[0].add_run()
    t3_run.text = "Next Steps"
    t3_run.font.name = "Calibri"
    t3_run.font.size = Pt(32)
    t3_run.font.bold = False

    # Body content
    body_tf3 = slide3.placeholders[1].text_frame
    body_tf3.paragraphs[0].text = ""

    next_items = [
        "Finalize partnership agreement with Meridian Corp by Oct 15",
        "Launch updated pricing tiers for mid-market segment",
        "Schedule onboarding sessions for 12 new enterprise clients",
        "Prepare board presentation with updated financial projections",
    ]
    for i, item in enumerate(next_items):
        if i == 0:
            p = body_tf3.paragraphs[0]
        else:
            p = body_tf3.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(18)
        run.font.bold = False

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
