"""
Reward Script: Set all title fonts to Century Gothic at 36pt
Task ID: osworld_impress_global_font_change_015
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5 pts): All 7 title placeholders use font name 'Century Gothic'
                          Partial credit: each correct title = 0.5/7 pts
  Component 2 (0.5 pts): All 7 title placeholders have font size 36pt (457200 EMU)
                          Partial credit: each correct title = 0.5/7 pts
Total: 1.0

Context:
  Initial state: All 7 title placeholders use Calibri at 40pt (508000 EMU)
  Golden state:  All 7 title placeholders use Century Gothic at 36pt (457200 EMU)
  Body/content text (ph_idx != 0) is NOT changed and NOT scored.
"""

import os

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError as e:
    print(f"CRITICAL: python-pptx not available: {e}")
    print("REWARD: 0.0")
    exit(1)

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_015'

EXPECTED_FONT_NAME = 'Century Gothic'
EXPECTED_FONT_SIZE_EMU = 457200  # 36pt = 36 * 12700 = 457200 EMU
EXPECTED_SLIDE_COUNT = 7


def persist_app_state():
    """Try to save any open LibreOffice presentation before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that all 7 title placeholders have been changed to Century Gothic at 36pt.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: verify we have the expected 7 slides
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDE_COUNT:
        print(f"WARN: Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}. Continuing verification.")

    # Collect all title placeholders (ph_idx == 0) across all slides
    title_shapes = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if (hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                    and shape.has_text_frame):
                title_shapes.append((i + 1, shape))

    total_titles = len(title_shapes)
    if total_titles == 0:
        print("CRITICAL: No title placeholders (ph_idx=0) found in presentation.")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {total_titles} title placeholder(s) across {num_slides} slides.")

    # Component 1: All title placeholders use font name 'Century Gothic' (0.5 points)
    # Partial credit: each correct title = 0.5 / total_titles
    font_name_correct = 0
    font_name_details = []
    try:
        for slide_num, shape in title_shapes:
            shape_font_name = None
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        shape_font_name = run.font.name
                        break
                if shape_font_name is not None:
                    break

            if shape_font_name == EXPECTED_FONT_NAME:
                font_name_correct += 1
                font_name_details.append(f"  Slide {slide_num}: PASS (font={shape_font_name})")
            else:
                font_name_details.append(
                    f"  Slide {slide_num}: FAIL (expected '{EXPECTED_FONT_NAME}', found '{shape_font_name}')"
                )

        comp1_score = 0.5 * (font_name_correct / total_titles)
        for detail in font_name_details:
            print(detail)

        if font_name_correct > 0:
            total_score += comp1_score
            print(f"PASS/PARTIAL: Component 1 — {font_name_correct}/{total_titles} title(s) use '{EXPECTED_FONT_NAME}' ({comp1_score:.4f} pts)")
        else:
            print(f"FAIL: Component 1 — 0/{total_titles} title(s) use '{EXPECTED_FONT_NAME}' (0.0 pts)")

    except Exception as e:
        print(f"ERROR: Component 1 (font name check) failed: {e}")

    # Component 2: All title placeholders have font size 36pt (457200 EMU) (0.5 points)
    # Partial credit: each correct title = 0.5 / total_titles
    font_size_correct = 0
    font_size_details = []
    try:
        for slide_num, shape in title_shapes:
            shape_font_size = None
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        shape_font_size = run.font.size
                        break
                if shape_font_size is not None:
                    break

            if shape_font_size == EXPECTED_FONT_SIZE_EMU:
                font_size_correct += 1
                font_size_details.append(
                    f"  Slide {slide_num}: PASS (size={shape_font_size} EMU = {shape_font_size / 12700:.1f}pt)"
                )
            else:
                actual_pt = (shape_font_size / 12700) if shape_font_size else None
                font_size_details.append(
                    f"  Slide {slide_num}: FAIL (expected {EXPECTED_FONT_SIZE_EMU} EMU = 36pt, "
                    f"found {shape_font_size} EMU = {actual_pt}pt)"
                )

        comp2_score = 0.5 * (font_size_correct / total_titles)
        for detail in font_size_details:
            print(detail)

        if font_size_correct > 0:
            total_score += comp2_score
            print(f"PASS/PARTIAL: Component 2 — {font_size_correct}/{total_titles} title(s) have size 36pt ({comp2_score:.4f} pts)")
        else:
            print(f"FAIL: Component 2 — 0/{total_titles} title(s) have size 36pt (0.0 pts)")

    except Exception as e:
        print(f"ERROR: Component 2 (font size check) failed: {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
