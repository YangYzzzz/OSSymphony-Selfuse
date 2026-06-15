"""
Initial Setup: 7-slide marketing pitch deck with Calibri 40pt titles
Task ID: osworld_impress_global_font_change_015
Domain: libreoffice_impress

Creates a presentation where all title placeholders use Calibri at 40pt.
The task is to change all titles to Century Gothic at 36pt.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_global_font_change_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_font(title_tf, text, font_name="Calibri", font_size_pt=40):
    """Set title placeholder text and font."""
    title_tf.text = text
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)


def create_initial():
    prs = Presentation()

    # Slide dimensions: standard widescreen 13.33 x 7.5 inches
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide data: realistic marketing pitch deck content ---
    slides_data = [
        {
            "layout": 0,  # Title Slide
            "title": "NovaSpark Solutions",
            "subtitle": "Transforming Digital Engagement for Modern Enterprises",
        },
        {
            "layout": 1,  # Title + Content
            "title": "Market Opportunity",
            "content": (
                "Global digital transformation market valued at $2.3 trillion\n"
                "Annual growth rate of 23% through 2028\n"
                "SMBs underserved by current enterprise solutions\n"
                "80% of businesses plan to increase tech spending in next 2 years"
            ),
        },
        {
            "layout": 1,
            "title": "Our Product Suite",
            "content": (
                "DataBridge Analytics Platform\n"
                "CloudSync Collaboration Tools\n"
                "SecureVault Identity Management\n"
                "InsightPulse Customer Intelligence"
            ),
        },
        {
            "layout": 1,
            "title": "Competitive Advantage",
            "content": (
                "Proprietary AI engine with 94% accuracy\n"
                "On-premise and hybrid cloud deployment\n"
                "Integration with 200+ enterprise tools\n"
                "99.98% uptime SLA guaranteed"
            ),
        },
        {
            "layout": 1,
            "title": "Traction & Milestones",
            "content": (
                "150+ enterprise clients in 18 months\n"
                "ARR grew from $1.2M to $8.7M (Q1 2024 – Q1 2025)\n"
                "NPS score of 72 across customer base\n"
                "Expansion into APAC and EMEA markets"
            ),
        },
        {
            "layout": 1,
            "title": "Financial Projections",
            "content": (
                "FY2025 Revenue Target: $18M\n"
                "Projected EBITDA margin: 22% by FY2026\n"
                "Customer acquisition cost: $4,200 (down 35% YoY)\n"
                "Lifetime value: $63,000 average"
            ),
        },
        {
            "layout": 0,  # Title Slide (closing)
            "title": "Join Us in Shaping the Future",
            "subtitle": "Contact: partnerships@novaspark.io  |  www.novaspark.io",
        },
    ]

    for idx, slide_info in enumerate(slides_data):
        layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(layout)

        # Set title — always Calibri at 40pt (initial state)
        if slide.shapes.title is not None:
            set_title_font(
                slide.shapes.title.text_frame,
                slide_info["title"],
                font_name="Calibri",
                font_size_pt=40,
            )

        # Set body / subtitle content
        if slide_info["layout"] == 0:
            # Title Slide: placeholders[1] = subtitle
            try:
                ph = slide.placeholders[1]
                ph.text = slide_info.get("subtitle", "")
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(20)
            except (KeyError, IndexError):
                pass
        else:
            # Title+Content: placeholders[1] = content
            try:
                ph = slide.placeholders[1]
                tf = ph.text_frame
                tf.text = ""
                lines = slide_info.get("content", "").split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = line
                    para.level = 0
                    for run in para.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(18)
            except (KeyError, IndexError):
                pass

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
