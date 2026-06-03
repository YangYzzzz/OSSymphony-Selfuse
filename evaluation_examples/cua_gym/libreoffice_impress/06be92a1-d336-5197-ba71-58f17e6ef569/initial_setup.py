"""
Initial Setup: Semester review presentation for Outline view reorganization
Task ID: impress_gf5_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, title_text, bullets, level=0):
    """Add a slide with title and bullet content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = level
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Semester Review - Fall 2025"
    slide1.placeholders[1].text = "Department of Computer Science\nProfessor Elena Rodriguez"

    # --- Slide 2: Course Overview ---
    add_title_content_slide(prs, "Course Overview", [
        "Enrolled 287 students across 4 undergraduate courses",
        "Added new elective: Machine Learning Fundamentals (CS 4820)",
        "Average class size increased by 12% compared to Spring 2025",
        "Introduced weekly coding labs for CS 2100 and CS 3300",
        "Guest lectures from industry partners: Anthropic, Stripe, Databricks",
    ])

    # --- Slide 3: Teaching Methodology ---
    add_title_content_slide(prs, "Teaching Methodology", [
        "Shifted to flipped classroom model for CS 3300 Algorithms",
        "Piloted AI-assisted grading for homework submissions",
        "Implemented peer code review sessions in all upper-division courses",
        "Adopted new textbook for CS 2100 Data Structures (Skiena, 4th ed.)",
    ])

    # --- Slide 4: Student Performance ---
    add_title_content_slide(prs, "Student Performance", [
        "Overall GPA improved from 3.12 to 3.28 across department courses",
        "Pass rate for CS 2100: 91% (up from 84% in Spring 2025)",
        "Senior capstone projects: 23 completed, 5 selected for conference",
        "Three students received NSF Graduate Research Fellowships",
    ])

    # --- Slide 5: Research Highlights (5 bullets, bullets 2-4 to be promoted) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Research Highlights"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.clear()

    bullets_5 = [
        ("Published 14 peer-reviewed papers in top-tier venues", 0),
        ("Natural Language Processing Lab Achievements", 1),
        ("Computer Vision Group Milestones", 1),
        ("Systems Security Research Breakthroughs", 1),
        ("Secured $2.3M in combined federal and industry grants", 0),
    ]
    for i, (text, level) in enumerate(bullets_5):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = text
        p.level = level
        for run in p.runs:
            run.font.size = Pt(18)

    # --- Slide 6: Department Initiatives ---
    add_title_content_slide(prs, "Department Initiatives", [
        "Launched mentorship program pairing seniors with freshmen",
        "Expanded internship partnerships to 18 companies",
        "Renovated Lab 204 with new GPU workstations",
        "Hosted first annual Hackathon with 150 participants",
    ])

    # --- Slide 7: Budget Summary ---
    add_title_content_slide(prs, "Budget Summary", [
        "Total departmental budget: $4.7M (5% increase from FY2024)",
        "Equipment and infrastructure: $890,000 allocated",
        "Travel and conference funding: $125,000 distributed among 32 faculty",
        "Remaining discretionary funds: $47,200 carried to Spring 2026",
    ])

    # --- Slide 8: Future Plans ---
    add_title_content_slide(prs, "Future Plans", [
        "Propose new M.S. track in Applied AI starting Fall 2026",
        "Hire 3 additional tenure-track faculty in emerging areas",
        "Expand online course offerings to reach non-traditional students",
        "Establish industry advisory board with quarterly meetings",
    ])

    # --- Slide 9: Acknowledgments ---
    add_title_content_slide(prs, "Acknowledgments", [
        "Dean Patricia Huang for unwavering support and advocacy",
        "IT Services team for seamless infrastructure upgrades",
        "Industry partners: Anthropic, Stripe, Databricks, NVIDIA",
        "Graduate TAs who went above and beyond this semester",
        "Administrative staff: Maria Santos and James O'Brien",
    ])

    # --- Slide 10: Questions & Discussion ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Questions & Discussion"
    body10 = slide10.placeholders[1]
    tf10 = body10.text_frame
    tf10.clear()
    p = tf10.paragraphs[0]
    p.text = "Thank you for your attention."
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(24)

    p2 = tf10.add_paragraph()
    p2.text = "Contact: e.rodriguez@university.edu"
    p2.alignment = PP_ALIGN.CENTER
    for run in p2.runs:
        run.font.size = Pt(18)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
