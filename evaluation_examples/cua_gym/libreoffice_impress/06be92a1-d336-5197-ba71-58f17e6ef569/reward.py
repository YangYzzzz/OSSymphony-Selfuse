"""
Reward Script: Reorganize semester_review.pptx outline structure
Task ID: impress_gf5_020
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide count is 12 (was 10, +3 promoted, -1 merged)
  Component 2 (0.30): Three promoted bullet points now exist as separate slide titles
  Component 3 (0.20): Original slide 5 no longer has the 3 level-1 bullets
  Component 4 (0.25): Slide with "Future Plans" title has Acknowledgments content in notes
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_020'

# The three bullet texts that should be promoted from level-1 to slide titles
PROMOTED_TITLES = [
    "Natural Language Processing Lab Achievements",
    "Computer Vision Group Milestones",
    "Systems Security Research Breakthroughs",
]

# Acknowledgments content that should appear in notes of the "Future Plans" slide
ACKNOWLEDGMENT_FRAGMENTS = [
    "Dean Patricia Huang",
    "IT Services team",
    "Anthropic, Stripe, Databricks, NVIDIA",
    "Graduate TAs",
    "Maria Santos",
]


def get_all_slide_titles(prs):
    """Extract the title text for each slide (first shape with title-like text)."""
    titles = []
    for slide in prs.slides:
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Take the first non-empty paragraph as the title candidate
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        title_text = para.text.strip()
                        break
            if title_text:
                break
        titles.append(title_text)
    return titles


def get_slide_notes(slide):
    """Safely get notes text from a slide."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    titles = get_all_slide_titles(prs)

    # Component 1: Slide count is 12 (0.25 points)
    # Initial has 10. After promoting 3 bullets to slides and deleting old slide 9: 10 + 3 - 1 = 12
    try:
        if num_slides == 12:
            print(f"PASS: Component 1 -- Slide count is 12 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Expected 12 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Three promoted bullets exist as separate slide titles (0.30 points)
    # Each promoted title is worth 0.10 points
    try:
        promoted_found = 0
        for expected_title in PROMOTED_TITLES:
            found_as_title = any(expected_title.lower() in t.lower() for t in titles)
            if found_as_title:
                promoted_found += 1
                print(f"PASS: Component 2 -- '{expected_title}' found as slide title (0.10 pts)")
            else:
                print(f"FAIL: Component 2 -- '{expected_title}' not found as any slide title")
        if promoted_found > 0:
            total_score += 0.10 * promoted_found
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: The "Research Highlights" slide no longer has level-1 bullets for the promoted items (0.20 points)
    # In the initial file, slide 5 had those as level-1 bullets. After promotion they should be gone.
    try:
        research_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if "Research Highlights" in para.text:
                            research_slide = slide
                            break
                if research_slide:
                    break
            if research_slide:
                break

        if research_slide is not None:
            # Collect all text from this slide
            all_text_on_slide = []
            for shape in research_slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_text_on_slide.append(para.text.strip())

            promoted_still_present = 0
            for pt in PROMOTED_TITLES:
                if any(pt.lower() in t.lower() for t in all_text_on_slide):
                    promoted_still_present += 1

            if promoted_still_present == 0:
                print(f"PASS: Component 3 -- Promoted bullets removed from Research Highlights slide (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 3 -- {promoted_still_present}/3 promoted bullets still on Research Highlights slide")
        else:
            print(f"FAIL: Component 3 -- 'Research Highlights' slide not found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: "Future Plans" slide has Acknowledgments content in speaker notes (0.25 points)
    # The old slide 9 (Acknowledgments) content should be pasted into the notes of the slide titled "Future Plans"
    try:
        future_plans_slide = None
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if "Future Plans" in para.text:
                            future_plans_slide = slide
                            break
                if future_plans_slide:
                    break
            if future_plans_slide:
                break

        if future_plans_slide is not None:
            notes_text = get_slide_notes(future_plans_slide)
            if notes_text:
                # Check how many acknowledgment fragments appear in the notes
                frags_found = 0
                for frag in ACKNOWLEDGMENT_FRAGMENTS:
                    if frag.lower() in notes_text.lower():
                        frags_found += 1

                if frags_found >= 3:
                    print(f"PASS: Component 4 -- Acknowledgments content found in Future Plans notes ({frags_found}/{len(ACKNOWLEDGMENT_FRAGMENTS)} fragments) (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 4 -- Only {frags_found}/{len(ACKNOWLEDGMENT_FRAGMENTS)} acknowledgment fragments in Future Plans notes")
                    print(f"  Notes content: {notes_text[:200]}")
            else:
                print(f"FAIL: Component 4 -- Future Plans slide has no speaker notes")
        else:
            print(f"FAIL: Component 4 -- 'Future Plans' slide not found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
