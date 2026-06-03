"""
Initial Setup: Create a 4-slide presentation with a table on slide 3.
Task ID: impress_tct_014
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Schedule Plan"
    slide1.placeholders[1].text = "Marketing Department — July to September 2025"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Overview"
    content = slide2.placeholders[1]
    content.text = "This quarter focuses on three major initiatives:"
    tf = content.text_frame
    for item in [
        "Brand refresh campaign launch across digital channels",
        "Customer engagement analytics dashboard rollout",
        "Partnership program expansion with key vendors",
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- Slide 3: Schedule Table ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title text box
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Detailed Weekly Schedule"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Table: 7 rows x 4 columns, 9 inches wide
    # All columns equal width: 2.25 inches each
    table_width = Inches(9)
    table_height = Inches(4.5)
    table_left = Inches(2.0)
    table_top = Inches(1.3)

    table_shape = slide3.shapes.add_table(
        7, 4, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Set all columns to equal width (2.25 inches each)
    col_width = Inches(2.25)
    for col_idx in range(4):
        table.columns[col_idx].width = col_width

    # Headers
    headers = ["Week", "Task", "Owner", "Status"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Dark header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x6E)

    # Table data
    data = [
        ["Week 1 (Jul 7)", "Finalize brand assets", "Sarah Chen", "Pending"],
        ["Week 2 (Jul 14)", "Launch social media ads", "Marcus Johnson", "Pending"],
        ["Week 3 (Jul 21)", "Analytics dashboard beta", "Priya Patel", "In Progress"],
        ["Week 4 (Jul 28)", "Partner outreach calls", "David Kim", "Scheduled"],
        ["Week 5 (Aug 4)", "Campaign performance review", "Sarah Chen", "Not Started"],
        ["Week 6 (Aug 11)", "Dashboard user testing", "Priya Patel", "Not Started"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
            if c == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Slide 4: Summary ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Milestones"
    content4 = slide4.placeholders[1]
    content4.text = "Critical dates for Q3 delivery:"
    tf4 = content4.text_frame
    milestones = [
        "July 15 — Brand campaign go-live",
        "August 1 — Dashboard v1.0 release",
        "August 20 — Partner agreements signed",
        "September 5 — Quarter review meeting",
    ]
    for m in milestones:
        p = tf4.add_paragraph()
        p.text = m
        p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
