"""
Reward Script: Set the first column width to 2.5 inches and make all remaining columns equal width
Task ID: impress_tct_014
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): First column width is 2.5 inches
  Component 2 (0.5): Remaining 3 columns are equal width (~2.1667 inches each) AND total width preserved at 9 inches
"""

import os
import time

from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_014'

# Tolerance for EMU comparisons (0.5% relative)
def is_approx_equal(val1, val2, tolerance=0.005):
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return val1 == val2
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    # Find table on slide 3 (index 2)
    slide = prs.slides[2]
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            table_shape = shape
            break

    if table is None:
        print("FAIL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: table is 4 columns x 7 rows
    num_cols = len(table.columns)
    num_rows = len(table.rows)
    if num_cols != 4 or num_rows != 7:
        print(f"FAIL: Table dimensions are {num_rows}x{num_cols}, expected 7x4")
        print("REWARD: 0.0")
        return 0.0

    # Get column widths
    col_widths = [table.columns[i].width for i in range(num_cols)]
    print(f"INFO: Column widths (EMU): {col_widths}")
    print(f"INFO: Column widths (inches): {[w/914400 for w in col_widths]}")

    # Expected values
    expected_col0_width = Inches(2.5)  # 2286000 EMU
    expected_remaining_width = (Inches(9) - expected_col0_width) / 3  # ~1981200 EMU each

    # Component 1: First column width is 2.5 inches (0.5 points)
    # This FAILS on initial (2.25 in) and PASSES on golden (2.5 in)
    try:
        col0_width = col_widths[0]
        if is_approx_equal(col0_width, expected_col0_width):
            print(f"PASS: Component 1 -- First column width is {col0_width/914400:.4f} inches (~2.5 inches) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 -- First column width is {col0_width/914400:.4f} inches, expected ~2.5 inches")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Remaining 3 columns are equal width AND total table width preserved (0.5 points)
    # This FAILS on initial (all cols 2.25 in, remaining are not ~2.1667) and PASSES on golden
    try:
        remaining_widths = col_widths[1:]
        total_width = sum(col_widths)
        expected_total = Inches(9)  # 8229600 EMU

        # Sub-check A: all remaining columns are equal to each other
        all_equal = all(is_approx_equal(w, remaining_widths[0]) for w in remaining_widths)
        # Sub-check B: remaining columns are the correct width (~2.1667 inches)
        correct_width = is_approx_equal(remaining_widths[0], expected_remaining_width)
        # Sub-check C: total table width is still 9 inches
        total_preserved = is_approx_equal(total_width, expected_total)

        if all_equal and correct_width and total_preserved:
            avg_width = sum(remaining_widths) / len(remaining_widths)
            print(f"PASS: Component 2 -- Remaining columns equal at ~{avg_width/914400:.4f} in, total width {total_width/914400:.4f} in (0.5 pts)")
            total_score += 0.5
        else:
            if not all_equal:
                print(f"FAIL: Component 2 -- Remaining columns are not equal: {[w/914400 for w in remaining_widths]}")
            if not correct_width:
                print(f"FAIL: Component 2 -- Remaining column width {remaining_widths[0]/914400:.4f} in, expected ~{expected_remaining_width/914400:.4f} in")
            if not total_preserved:
                print(f"FAIL: Component 2 -- Total table width {total_width/914400:.4f} in, expected 9.0 in")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved state
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
