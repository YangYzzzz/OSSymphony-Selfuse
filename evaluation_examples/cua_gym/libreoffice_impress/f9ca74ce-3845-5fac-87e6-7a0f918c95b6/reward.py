"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m tidying up a huge deck in LibreOffice Impress and spotted an issue on slide 185: the first two bullet points really need to be numbered instead. What’s the quickest way to turn just those two bullets into a numbered list that starts at 1 without touching anything else on that slide?
Generated: 2025-09-10 18:23:44
Status: success
Model: azure-o3
Total Steps: 10
"""

import os
from pptx import Presentation

def _paragraph_type(paragraph):
    """Return the type of bulleting/numbering applied to a paragraph.
    Possible return values: 'number', 'bullet', 'none', 'inherit'"""
    auto  = paragraph._p.xpath('./a:pPr/a:buAutoNum')
    char  = paragraph._p.xpath('./a:pPr/a:buChar')
    none  = paragraph._p.xpath('./a:pPr/a:buNone')
    if auto:
        return 'number'
    if char:
        return 'bullet'
    if none:
        return 'none'
    return 'inherit'

def verify_slide_185_numbering(file_path):
    """Verify that on slide 185:   
    • the first two top-level paragraphs are numbered, starting at 1   
    • all subsequent paragraphs are NOT numbered.
    Progressive scoring (max 1.0):
        0.4 – correct numbering pattern detected (first two numbered, rest not)
        0.3 – first number starts at 1 (or startAt omitted which defaults to 1)
        0.3 – confirmation that paragraphs after the second are NOT numbered
    """
    score = 0.0
    max_score = 1.0

    print(f"Loading presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        return score

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error opening presentation: {e}")
        return score

    if len(prs.slides) < 185:
        print(f"✗ Expected at least 185 slides, found {len(prs.slides)}")
        return score

    slide = prs.slides[184]  # 0-based index → slide 185

    # Gather all non-empty, top-level paragraphs in reading order
    paragraphs = []
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for p in shape.text_frame.paragraphs:
            if p.text.strip():
                paragraphs.append(p)

    if len(paragraphs) < 3:
        print("✗ Not enough paragraphs found to perform verification")
        return score

    # Determine paragraph types
    types = [_paragraph_type(p) for p in paragraphs]

    # ------------- Scoring begins -------------
    # Condition A: first two numbered; all after second NOT numbered
    if types[0] == 'number' and types[1] == 'number' and all(t != 'number' for t in types[2:]):
        score += 0.4
        print("✓ Found expected numbering pattern (0.4)")

        # Condition B: first number starts at 1
        auto = paragraphs[0]._p.xpath('./a:pPr/a:buAutoNum')
        start_val = auto[0].get('startAt') if auto else None  # None ⇒ defaults to 1
        if start_val in (None, '1'):
            score += 0.3
            print("✓ First list item starts at 1 (0.3)")
        else:
            print(f"✗ First list item should start at 1, found startAt={start_val}")

        # Condition C: ensure every paragraph after the 2nd is NOT numbered
        if all(t != 'number' for t in types[2:]):
            score += 0.3
            print("✓ Subsequent paragraphs are not numbered (0.3)")
        else:
            print("✗ Some subsequent paragraphs are still numbered")
    else:
        print("✗ First two paragraphs are not uniquely numbered OR later paragraphs wrongly numbered")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score

# ----------------- MAIN EXECUTION -----------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/im_tidying_up_a_huge_deck_in_libreoffice_impress_and_spotted_an_issue_on_slide_185_the_first_two_bul_golden.pptx"
    reward = verify_slide_185_numbering(FILE_PATH)
    print(f"REWARD: {reward}")
