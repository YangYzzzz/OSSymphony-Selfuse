"""
Initial Setup: Create a data dashboard presentation with slide 8 as blank Gantt Chart placeholder
Task ID: impress_gf5_044
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_044'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_blank_titled_slide(prs, title_text):
    """Add a slide with only a title (layout 5 = blank, add textbox for title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5F)
    return slide


def add_table_slide(prs, title_text, headers, data):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3A, 0x5F)

    rows = len(data) + 1
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * rows))
    tbl = tbl_shape.table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = str(val)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Q2 2025 Data Dashboard", "Executive Strategy & Operations Review\nPrepared by Analytics Team")

    # Slide 2: Revenue Overview
    add_content_slide(prs, "Revenue Overview", [
        "Total Revenue: $4.82M (+12.3% YoY)",
        "Recurring Revenue: $3.15M (65.4% of total)",
        "New Customer Revenue: $1.67M",
        "Average Deal Size: $48,200",
        "Pipeline Value: $7.3M",
    ])

    # Slide 3: Regional Performance Table
    add_table_slide(prs, "Regional Performance Summary", [
        "Region", "Revenue", "Growth", "Customers", "Avg Deal"
    ], [
        ["North America", "$2.14M", "+15.2%", "89", "$54,300"],
        ["Europe", "$1.38M", "+9.7%", "62", "$44,800"],
        ["Asia Pacific", "$0.87M", "+18.4%", "41", "$38,200"],
        ["Latin America", "$0.43M", "+6.1%", "23", "$31,500"],
    ])

    # Slide 4: Product Metrics
    add_content_slide(prs, "Product Metrics", [
        "Monthly Active Users: 142,800 (+22% QoQ)",
        "Feature Adoption Rate: 67.3%",
        "NPS Score: 72 (up from 65)",
        "Avg Session Duration: 14.2 minutes",
        "Support Tickets: 1,240 (-8% MoM)",
        "Uptime: 99.97%",
    ])

    # Slide 5: Team Headcount
    add_table_slide(prs, "Team Headcount by Department", [
        "Department", "Current", "Open Roles", "Q3 Target", "Budget"
    ], [
        ["Engineering", "48", "6", "54", "$890K"],
        ["Product", "12", "2", "14", "$245K"],
        ["Sales", "34", "8", "42", "$620K"],
        ["Marketing", "18", "3", "21", "$380K"],
        ["Operations", "15", "1", "16", "$210K"],
        ["Customer Success", "22", "4", "26", "$345K"],
    ])

    # Slide 6: Key Initiatives
    add_content_slide(prs, "Key Strategic Initiatives", [
        "1. Enterprise Tier Launch - Target: June 15",
        "2. Mobile App v2.0 - Beta testing underway",
        "3. APAC Market Expansion - Singapore office setup",
        "4. AI-Powered Analytics Module - Phase 2 development",
        "5. SOC 2 Type II Certification - Audit in progress",
    ])

    # Slide 7: Customer Satisfaction
    add_table_slide(prs, "Customer Satisfaction Trends", [
        "Metric", "Q1 2025", "Q2 2025", "Target", "Status"
    ], [
        ["NPS Score", "65", "72", "75", "On Track"],
        ["CSAT Rating", "4.2", "4.5", "4.6", "On Track"],
        ["Churn Rate", "3.8%", "3.1%", "2.5%", "Improving"],
        ["Renewal Rate", "87%", "91%", "93%", "On Track"],
        ["Response Time", "4.2h", "2.8h", "2.0h", "Improving"],
    ])

    # Slide 8: Project Gantt Chart - BLANK with title only (agent must create the Gantt chart)
    add_blank_titled_slide(prs, "Project Gantt Chart")

    # Slide 9: Risk Register
    add_table_slide(prs, "Risk Register", [
        "Risk", "Severity", "Likelihood", "Mitigation"
    ], [
        ["Supply chain delays", "High", "Medium", "Dual-source strategy"],
        ["Key talent attrition", "High", "Low", "Retention packages"],
        ["Regulatory changes", "Medium", "Medium", "Compliance monitoring"],
        ["Competitor pricing war", "Medium", "High", "Value differentiation"],
    ])

    # Slide 10: Next Steps
    add_content_slide(prs, "Next Steps & Action Items", [
        "Complete enterprise pricing model by May 30",
        "Finalize APAC hiring plan with HR by June 5",
        "Submit SOC 2 documentation by June 10",
        "Launch beta invitations for mobile v2.0 by June 12",
        "Board presentation preparation due June 20",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
