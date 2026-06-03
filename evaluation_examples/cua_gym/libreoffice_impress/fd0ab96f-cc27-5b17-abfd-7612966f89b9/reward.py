"""
Reward Script: Gantt chart creation on slide 8 using drawing shapes
Task ID: impress_gf5_044
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Slide 8 has significantly more shapes than initial (>=15 new shapes)
  Component 2 (0.25): 8 task bars exist as colored rectangles/rounded rectangles
  Component 3 (0.20): Color coding by phase (Planning=blue, Dev=green, Testing=orange, Launch=red)
  Component 4 (0.15): Date axis labels (Jan-Jun) present at top of slide
  Component 5 (0.20): Each task bar has a text label with the task name
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_044'

# Expected task names for the 8 tasks
EXPECTED_TASKS = [
    'Requirements Gathering',
    'Architecture Design',
    'Backend Development',
    'Frontend Development',
    'API Integration',
    'QA Testing',
    'User Acceptance Testing',
    'Product Launch',
]

# Phase color mapping (approximate RGB values)
PHASE_COLORS = {
    'Planning': {'r_range': (0, 120), 'g_range': (0, 160), 'b_range': (150, 255)},   # blue-ish
    'Development': {'r_range': (0, 100), 'g_range': (100, 255), 'b_range': (0, 130)}, # green-ish
    'Testing': {'r_range': (180, 255), 'g_range': (90, 180), 'b_range': (0, 100)},    # orange-ish
    'Launch': {'r_range': (150, 255), 'g_range': (0, 100), 'b_range': (0, 100)},      # red-ish
}

# Map tasks to phases
TASK_PHASES = {
    'Requirements Gathering': 'Planning',
    'Architecture Design': 'Planning',
    'Backend Development': 'Development',
    'Frontend Development': 'Development',
    'API Integration': 'Development',
    'QA Testing': 'Testing',
    'User Acceptance Testing': 'Testing',
    'Product Launch': 'Launch',
}

MONTH_LABELS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']


def color_matches_phase(rgb, phase):
    """Check if an RGB color roughly matches the expected phase color."""
    if rgb is None:
        return False
    r, g, b = rgb[0], rgb[1], rgb[2]
    pr = PHASE_COLORS[phase]
    return (pr['r_range'][0] <= r <= pr['r_range'][1] and
            pr['g_range'][0] <= g <= pr['g_range'][1] and
            pr['b_range'][0] <= b <= pr['b_range'][1])


def get_shape_fill_rgb(shape):
    """Try to get the fill color RGB of a shape."""
    try:
        fill = shape.fill
        if fill.type is not None:
            rgb = fill.fore_color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None


def get_all_text_from_shape(shape):
    """Get all text from a shape including text frame."""
    texts = []
    if hasattr(shape, 'text_frame'):
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    elif hasattr(shape, 'text') and shape.text:
        texts.append(shape.text.strip())
    return texts


def verify_task(file_path):
    """
    Verify Gantt chart creation with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Check we have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # slide 8, 0-indexed
    shapes = list(slide.shapes)
    num_shapes = len(shapes)

    # Collect all text on the slide
    all_texts = []
    for shape in shapes:
        all_texts.extend(get_all_text_from_shape(shape))

    # Collect colored auto shapes (potential task bars)
    colored_bars = []
    for shape in shapes:
        if shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, 1):
            rgb = get_shape_fill_rgb(shape)
            if rgb is not None:
                # Filter out very thin shapes (gridlines) and very small shapes (legend squares)
                # Task bars should be wider than they are tall, and reasonably sized
                w = shape.width
                h = shape.height
                # Bars: width > 300000 EMU (~0.33 inch) and height > 200000 EMU
                if w > 300000 and h > 200000:
                    text = shape.text.strip() if hasattr(shape, 'text') and shape.text else ''
                    colored_bars.append({
                        'name': shape.name,
                        'text': text,
                        'rgb': rgb,
                        'left': shape.left,
                        'top': shape.top,
                        'width': w,
                        'height': h,
                    })

    print(f"Slide 8: {num_shapes} total shapes, {len(colored_bars)} colored bars detected")

    # ---- Component 1: Slide 8 has significantly more shapes (>= 15 new) (0.20 pts) ----
    # Initial slide 8 has 2 shapes. Golden should have many more.
    try:
        if num_shapes >= 17:  # at least 15 new shapes added
            print(f"PASS: Component 1 - {num_shapes} shapes on slide 8 (>=17 required) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 - Only {num_shapes} shapes on slide 8 (need >=17)")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # ---- Component 2: 8 task bars exist as colored rectangles (0.25 pts) ----
    try:
        if len(colored_bars) >= 8:
            print(f"PASS: Component 2 - {len(colored_bars)} task bars found (>=8 required) (0.25 pts)")
            total_score += 0.25
        elif len(colored_bars) >= 5:
            partial = 0.15
            print(f"PARTIAL: Component 2 - {len(colored_bars)} task bars found (5-7), awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - Only {len(colored_bars)} task bars found (need >=8)")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # ---- Component 3: Color coding by phase (0.20 pts) ----
    try:
        # Check that bars have at least 2 distinct color families
        # Group bars by approximate color family
        color_families = set()
        for bar in colored_bars:
            r, g, b = bar['rgb']
            for phase_name, _ in PHASE_COLORS.items():
                if color_matches_phase(bar['rgb'], phase_name):
                    color_families.add(phase_name)
                    break

        num_families = len(color_families)
        print(f"  Color families detected: {color_families}")

        if num_families >= 4:
            print(f"PASS: Component 3 - All 4 phase colors present (0.20 pts)")
            total_score += 0.20
        elif num_families >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 3 - {num_families} phase colors present, awarding {partial} pts")
            total_score += partial
        elif num_families >= 2:
            partial = 0.10
            print(f"PARTIAL: Component 3 - {num_families} phase colors present, awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - Only {num_families} distinct phase color(s)")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # ---- Component 4: Date axis labels (Jan-Jun) at top of slide (0.15 pts) ----
    try:
        found_months = []
        for month in MONTH_LABELS:
            for text in all_texts:
                if month.lower() in text.lower():
                    found_months.append(month)
                    break

        if len(found_months) >= 6:
            print(f"PASS: Component 4 - All 6 month labels found: {found_months} (0.15 pts)")
            total_score += 0.15
        elif len(found_months) >= 4:
            partial = 0.10
            print(f"PARTIAL: Component 4 - {len(found_months)} month labels found: {found_months}, awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - Only {len(found_months)} month labels found: {found_months}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # ---- Component 5: Task name labels present (0.20 pts) ----
    try:
        # Check how many of the 8 expected task names appear anywhere on the slide
        found_tasks = []
        for task_name in EXPECTED_TASKS:
            for text in all_texts:
                if task_name.lower() in text.lower():
                    found_tasks.append(task_name)
                    break

        if len(found_tasks) >= 8:
            print(f"PASS: Component 5 - All 8 task labels found (0.20 pts)")
            total_score += 0.20
        elif len(found_tasks) >= 5:
            partial = 0.12
            print(f"PARTIAL: Component 5 - {len(found_tasks)}/8 task labels found: {found_tasks}, awarding {partial} pts")
            total_score += partial
        elif len(found_tasks) >= 3:
            partial = 0.06
            print(f"PARTIAL: Component 5 - {len(found_tasks)}/8 task labels found: {found_tasks}, awarding {partial} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 5 - Only {len(found_tasks)}/8 task labels found: {found_tasks}")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
