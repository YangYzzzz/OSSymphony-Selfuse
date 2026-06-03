"""
Initial Setup: Create Quarterly_Review.pptx with 6 slides for reorganization task
Task ID: impress_wf_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Q3 Review (Title Slide) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Review"
    slide1.placeholders[1].text = "Quarterly Business Performance\nFiscal Year 2025"

    # --- Slide 2: Revenue ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total Q3 Revenue: $4.82M"
    p = body2.add_paragraph()
    p.text = "Product Sales: $2.95M (+12% YoY)"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Service Revenue: $1.47M (+8% YoY)"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Licensing Fees: $400K (+3% YoY)"
    p.level = 1
    p = body2.add_paragraph()
    p.text = ""
    p = body2.add_paragraph()
    p.text = "Revenue exceeded projections by 6.2% this quarter."

    # --- Slide 3: Expenses ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Expenses"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Q3 Expenses: $3.15M"
    p = body3.add_paragraph()
    p.text = "Personnel Costs: $1.82M (57.8%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Infrastructure & Cloud: $620K (19.7%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Marketing & Sales: $410K (13.0%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "R&D Investment: $300K (9.5%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = ""
    p = body3.add_paragraph()
    p.text = "Operating margin improved to 34.6%, up from 31.2% in Q2."

    # --- Slide 4: Deprecated Info ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Deprecated Info"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Legacy System Migration Status"
    p = body4.add_paragraph()
    p.text = "Oracle DB decommission: Complete"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "On-prem servers retired: 14 of 14"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Data migration validation: Passed"
    p.level = 1
    p = body4.add_paragraph()
    p.text = ""
    p = body4.add_paragraph()
    p.text = "This slide is no longer relevant for quarterly reviews."

    # --- Slide 5: Key Highlights ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Highlights"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Q3 2025 Achievements"
    p = body5.add_paragraph()
    p.text = "Launched Enterprise Pro tier with 340 signups in first month"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Customer NPS increased from 62 to 71"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Closed strategic partnership with Meridian Corp ($1.2M ARR)"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Engineering team shipped 94% of planned sprint objectives"
    p.level = 1
    p = body5.add_paragraph()
    p.text = "Employee retention rate: 96.3%"
    p.level = 1

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Q4 2025 Priorities"
    p = body6.add_paragraph()
    p.text = "Expand APAC sales team by 5 headcount"
    p.level = 1
    p = body6.add_paragraph()
    p.text = "Launch mobile app v2.0 by November 15"
    p.level = 1
    p = body6.add_paragraph()
    p.text = "Achieve SOC 2 Type II certification"
    p.level = 1
    p = body6.add_paragraph()
    p.text = "Target $5.5M revenue for Q4"
    p.level = 1
    p = body6.add_paragraph()
    p.text = "Conduct annual customer satisfaction survey"
    p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
