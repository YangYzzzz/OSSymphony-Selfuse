"""
Reward Script: Reorganize Quarterly_Review.pptx slides
Task ID: impress_wf_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): 'Deprecated Info' slide removed and slide count is 6
  Component 2 (0.3): Slides 1-5 in correct order
  Component 3 (0.3): Last slide title is 'Questions?'
  Component 4 (0.2): Last slide uses Title Slide layout (same as slide 1)
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_007'


def get_slide_title(slide):
    """Extract the first non-empty text from a slide's shapes."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)
    slide_titles = [get_slide_title(s) for s in slides]

    print(f"INFO: Found {num_slides} slides with titles: {slide_titles}")

    # Component 1: 'Deprecated Info' removed AND slide count is 6 (0.2 points)
    # Initial has 6 slides WITH 'Deprecated Info' -> this component FAILS on initial
    # Golden has 6 slides WITHOUT 'Deprecated Info' -> PASSES on golden
    try:
        has_deprecated = any("Deprecated Info" in t for t in slide_titles)
        if num_slides == 6 and not has_deprecated:
            print(f"PASS: Component 1 — 6 slides, no 'Deprecated Info' (0.2 pts)")
            total_score += 0.2
        else:
            if has_deprecated:
                print(f"FAIL: Component 1 — 'Deprecated Info' slide still present")
            if num_slides != 6:
                print(f"FAIL: Component 1 — Expected 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slides 1-5 in correct order (0.3 points)
    # Expected order: Q3 Review, Key Highlights, Revenue, Expenses, Next Steps
    # Initial order is: Q3 Review, Revenue, Expenses, Deprecated Info, Key Highlights, Next Steps
    # So this FAILS on initial (slide 2 should be 'Key Highlights' but is 'Revenue')
    try:
        expected_order = ['Q3 Review', 'Key Highlights', 'Revenue', 'Expenses', 'Next Steps']
        if num_slides >= 5:
            actual_first_five = slide_titles[:5]
            matches = sum(1 for a, e in zip(actual_first_five, expected_order) if a == e)
            if matches == 5:
                print(f"PASS: Component 2 — Slides 1-5 in correct order: {actual_first_five} (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Expected order {expected_order}, found {actual_first_five} ({matches}/5 match)")
        else:
            print(f"FAIL: Component 2 — Not enough slides ({num_slides}) to check order")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Last slide title is 'Questions?' (0.3 points)
    # Initial last slide is 'Next Steps' -> FAILS on initial
    try:
        if num_slides >= 1:
            last_title = slide_titles[-1]
            if last_title == 'Questions?':
                print(f"PASS: Component 3 — Last slide title is 'Questions?' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 3 — Expected last slide title 'Questions?', found '{last_title}'")
        else:
            print(f"FAIL: Component 3 — No slides found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Last slide uses Title Slide layout, matching slide 1 (0.2 points)
    # Initial last slide has 'Title and Content' layout -> FAILS on initial
    try:
        if num_slides >= 2:
            first_layout = slides[0].slide_layout.name
            last_layout = slides[-1].slide_layout.name
            if last_layout == first_layout:
                print(f"PASS: Component 4 — Last slide layout '{last_layout}' matches slide 1 layout (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 — Last slide layout '{last_layout}' does not match slide 1 layout '{first_layout}'")
        else:
            print(f"FAIL: Component 4 — Not enough slides to compare layouts")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
