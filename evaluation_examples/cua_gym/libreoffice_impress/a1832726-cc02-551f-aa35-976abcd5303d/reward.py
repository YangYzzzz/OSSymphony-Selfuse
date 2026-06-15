"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 13 I currently have a single image called “Picture 1.” I want to duplicate it (creating “Picture 2”) and then line both copies up side-by-side, making sure there’s exactly a 1.00 cm gap between their inside edges. How do I do that in LibreOffice Impress?
Generated: 2025-09-10 12:13:16
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Constants
EMU_PER_CM = 360000          # 1 cm in English Metric Units
GAP_TOLERANCE = 20000        # ±0.055 cm tolerance (≈0.55 mm)
VERTICAL_TOLERANCE = 20000   # Vertical alignment tolerance


def verify_picture_duplication(file_path: str) -> float:
    """Verify that on slide 13 of the given presentation:
    1. There are at least two pictures.
    2. They are specifically named “Picture 1” and “Picture 2”.
    3. They are aligned side-by-side (same top coordinate within tolerance).
    4. The gap between their inside edges is exactly 1.00 cm ± tolerance.

    Returns a progressive score from 0.0 to 1.0.
    """

    print(f"Loading presentation: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File does not exist")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # Load the presentation file
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # Ensure slide 13 exists (index 12, zero-based)
    target_slide_index = 12
    if len(prs.slides) <= target_slide_index:
        print(f"✗ Slide 13 not found (presentation has only {len(prs.slides)} slides)")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[target_slide_index]
    picture_shapes = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    print(f"Found {len(picture_shapes)} picture(s) on slide 13")

    total_score = 0.0  # progressive scoring
    # ------------------------------------------------------------------
    # Requirement 1: At least two pictures
    if len(picture_shapes) >= 2:
        total_score += 0.2
        print("✓ At least two pictures present (0.2)")
    else:
        print("✗ Less than two pictures – duplication not achieved")
        print(f"REWARD: {total_score}")
        return total_score  # cannot proceed with further checks

    # ------------------------------------------------------------------
    # Requirement 2: Correct names “Picture 1” and “Picture 2”
    names = [sh.name.strip().lower() for sh in picture_shapes]
    has_pic1 = "picture 1" in names
    has_pic2 = "picture 2" in names

    if has_pic1 and has_pic2:
        total_score += 0.2
        print("✓ Found shapes named ‘Picture 1’ and ‘Picture 2’ (0.2)")
    else:
        missing = [n for n, flag in [("Picture 1", has_pic1), ("Picture 2", has_pic2)] if not flag]
        print(f"✗ Missing expected names: {', '.join(missing)}")

    # ------------------------------------------------------------------
    # Requirement 3 & 4: Alignment and exact 1 cm gap
    # Proceed only if at least two pictures exist
    if len(picture_shapes) >= 2:
        # Sort pictures by horizontal position (left attribute)
        pic_left, pic_right = sorted(picture_shapes, key=lambda sh: sh.left)[:2]

        # 3. Vertical alignment (top coordinates nearly equal)
        vertical_diff = abs(pic_left.top - pic_right.top)
        if vertical_diff <= VERTICAL_TOLERANCE:
            total_score += 0.3
            print(f"✓ Pictures vertically aligned within {vertical_diff} EMU (0.3)")
        else:
            print(f"✗ Vertical misalignment: {vertical_diff} EMU > {VERTICAL_TOLERANCE}")

        # 4. Gap between inside edges == 1 cm ± tolerance
        measured_gap = pic_right.left - (pic_left.left + pic_left.width)
        gap_diff = abs(measured_gap - EMU_PER_CM)
        print(f"Measured gap: {measured_gap} EMU, expected: {EMU_PER_CM} EMU, diff: {gap_diff} EMU")
        if gap_diff <= GAP_TOLERANCE:
            total_score += 0.3
            print("✓ Gap between pictures is 1.00 cm ± tolerance (0.3)")
        else:
            print("✗ Gap not within tolerance")

    # ------------------------------------------------------------------
    # Cap score at 1.0 and round for readability
    final_score = round(min(total_score, 1.0), 2)
    print(f"Total score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when run as a script
if __name__ == "__main__":
    presentation_path = "/home/user/on_slide_13_i_currently_have_a_single_image_called_picture_1_i_want_to_duplicate_it_creating_picture_golden.pptx"
    verify_picture_duplication(presentation_path)

