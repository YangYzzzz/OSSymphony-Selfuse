"""
Reward Script: Extract presenter notes from NeurIPS_Talk.pptx into neurips_notes.docx
Task ID: osworld_multi_apps_impress_notes_export_006
Domain: libreoffice_impress (multi_apps: pptx -> docx)
Scoring:
  Component 1: neurips_notes.docx exists and contains all 15 slide notes (0.5 pts)
  Component 2: Notes are in correct slide order — check specific phrases per slide (0.3 pts)
  Component 3: Notes are separated by blank lines and no slide-number metadata present (0.2 pts)
Total: 1.0
"""

import os

# python-docx for reading docx; python-pptx for reading pptx notes as ground truth
from docx import Document
from pptx import Presentation

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_006'

# Expected file paths on the VM
DOCX_PATH = os.path.join(WORKDIR, 'neurips_notes.docx')
PPTX_PATH = os.path.join(WORKDIR, 'NeurIPS_Talk.pptx')


def get_pptx_notes(pptx_path):
    """Extract notes from each slide of the pptx, returning a list of 15 strings."""
    prs = Presentation(pptx_path)
    notes_list = []
    for slide in prs.slides:
        try:
            text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            text = ''
        notes_list.append(text)
    return notes_list


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # ---- PRECONDITION: pptx source file must exist ----
    if not os.path.exists(PPTX_PATH):
        print(f"PRECONDITION FAIL: Source pptx not found: {PPTX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Load expected notes from pptx (ground truth from source file)
    try:
        expected_notes = get_pptx_notes(PPTX_PATH)
    except Exception as e:
        print(f"PRECONDITION FAIL: Cannot load pptx notes: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---- PRECONDITION: docx target file must exist ----
    if not os.path.exists(DOCX_PATH):
        print(f"FAIL: neurips_notes.docx does not exist at {DOCX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Load docx
    try:
        doc = Document(DOCX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load neurips_notes.docx: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Get non-empty paragraph texts from the docx
    paragraphs = [p.text for p in doc.paragraphs]
    non_empty_paragraphs = [p for p in paragraphs if p.strip()]

    print(f"INFO: docx has {len(paragraphs)} total paragraphs, {len(non_empty_paragraphs)} non-empty")
    print(f"INFO: pptx has {len(expected_notes)} slides with notes")

    # ---- Component 1: neurips_notes.docx exists and contains all 15 slide notes (0.5 pts) ----
    # Check: at least 15 non-empty content paragraphs exist (one per slide),
    # and all 15 notes are present (check first 60 chars of each note).
    try:
        notes_found = 0
        missing_notes = []
        for i, expected in enumerate(expected_notes):
            if not expected:
                # Skip slides with empty notes
                notes_found += 1
                continue
            # Check first ~60 chars of each note appear in the docx text
            first_60 = expected[:60].strip()
            found = any(first_60 in p for p in non_empty_paragraphs)
            if found:
                notes_found += 1
            else:
                missing_notes.append(f"Slide {i+1}: '{first_60[:40]}...'")

        if notes_found == len(expected_notes):
            print(f"PASS: Component 1 — All {len(expected_notes)} slide notes found in docx (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Only {notes_found}/{len(expected_notes)} slide notes found")
            for m in missing_notes:
                print(f"  Missing: {m}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---- Component 2: Notes are in correct slide order (0.3 pts) ----
    # Join all non-empty paragraphs and check that key phrases appear in order.
    # We check a representative phrase from each of the 15 slides.
    try:
        # Build a single text block from the docx (joining all paragraphs)
        docx_full_text = '\n'.join(paragraphs)

        # Key unique phrases from each slide (each phrase appears ONLY in that slide's notes)
        slide_key_phrases = [
            "Welcome everyone to this presentation",            # Slide 1 — unique opener
            "The fundamental question we set",                  # Slide 2 — unique opener
            "Key prior work includes Kaplan",                   # Slide 3 — unique opener
            "We trained 412 models spanning",                   # Slide 4 — unique opener
            "We parameterize loss as L(N,",                     # Slide 5 — unique opener
            "Our main result is the",                           # Slide 6 — unique opener
            "When measuring scaling on downstream",             # Slide 7 — unique opener
            "We conducted an ablation study",                   # Slide 8 — unique opener
            "A key finding of our",                             # Slide 9 — unique opener
            "When we include instruction tuning",               # Slide 10 — unique opener
            "Several important limitations must be",            # Slide 11 — unique opener
            "Our findings have several implications",           # Slide 12 — unique opener
            "Several directions are promising for",             # Slide 13 — unique opener
            "To summarize the main contributions",              # Slide 14 — unique opener
            "We thank the reviewers for",                       # Slide 15 — unique opener
        ]

        # Check all phrases appear in the docx
        phrases_found = []
        phrases_missing = []
        for i, phrase in enumerate(slide_key_phrases):
            if phrase in docx_full_text:
                phrases_found.append(i + 1)
            else:
                phrases_missing.append(f"Slide {i+1}: '{phrase}'")

        # Check order: positions of key phrases should be non-decreasing
        phrase_positions = []
        for phrase in slide_key_phrases:
            pos = docx_full_text.find(phrase)
            phrase_positions.append(pos)

        in_order = all(
            phrase_positions[i] < phrase_positions[i+1]
            for i in range(len(phrase_positions)-1)
            if phrase_positions[i] != -1 and phrase_positions[i+1] != -1
        )

        found_count = len(phrases_found)
        if found_count == 15 and in_order:
            print(f"PASS: Component 2 — All 15 slide-specific phrases found in correct order (0.3 pts)")
            total_score += 0.3
        elif found_count >= 12 and in_order:
            # Partial: most slides present and in order
            partial = round(0.3 * found_count / 15, 2)
            print(f"PARTIAL: Component 2 — {found_count}/15 phrases found in order, awarding {partial} pts")
            if partial > 0:
                total_score += partial
        else:
            print(f"FAIL: Component 2 — {found_count}/15 phrases found; in_order={in_order}")
            for m in phrases_missing:
                print(f"  Missing phrase: {m}")
            if not in_order:
                print(f"  Order issue: phrase positions: {phrase_positions}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---- Component 3: Blank line separators exist and no slide-number labels (0.2 pts) ----
    # Per task: "each slide's notes separated by a blank line", "no slide number labels or other metadata"
    try:
        # Check blank lines exist between notes blocks
        # The docx should have empty paragraphs (blank lines) between note texts
        empty_para_count = sum(1 for p in paragraphs if not p.strip())

        # With 15 slides and 1 blank line between each, there should be at least 14 empty paragraphs
        has_blank_separators = empty_para_count >= 14

        # Check no slide number labels (e.g., "Slide 1:", "Slide 1 -", "Slide 1.", "[Slide 1]")
        import re
        slide_label_pattern = re.compile(
            r'(slide\s+\d+\s*[:\-\.\)]|^\s*\d+\s*[:\-\.]\s)', re.IGNORECASE | re.MULTILINE
        )
        has_slide_labels = bool(slide_label_pattern.search(docx_full_text))

        separator_ok = has_blank_separators
        no_labels_ok = not has_slide_labels

        if separator_ok and no_labels_ok:
            print(f"PASS: Component 3 — Blank separators present ({empty_para_count} empty paras) and no slide labels (0.2 pts)")
            total_score += 0.2
        elif separator_ok and not no_labels_ok:
            print(f"PARTIAL: Component 3 — Blank separators present but slide labels found; awarding 0.1 pts")
            total_score += 0.1
        elif not separator_ok and no_labels_ok:
            print(f"PARTIAL: Component 3 — No slide labels but insufficient blank separators ({empty_para_count} found, need >=14); awarding 0.1 pts")
            total_score += 0.1
        else:
            print(f"FAIL: Component 3 — blank_separators={separator_ok} ({empty_para_count} empty paras), no_labels={no_labels_ok}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entrypoint: run verification
verify_task()
