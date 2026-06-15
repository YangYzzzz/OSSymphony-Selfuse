"""
Initial Setup: NeurIPS_Talk.pptx with detailed speaker notes, open in LibreOffice Impress
Task ID: osworld_multi_apps_impress_notes_export_006
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user/Desktop'
TASK_ID = 'NeurIPS_Talk'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    os.makedirs(WORKDIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, subtitle_or_content, notes_text)
    slides_data = [
        (
            "Scaling Laws for Neural Language Models: A Deep Dive",
            "NeurIPS 2024 | Research Track",
            "Welcome everyone to this presentation on scaling laws for neural language models. "
            "Today we'll be exploring how model performance scales with compute, data, and parameters. "
            "This work builds on the foundational Chinchilla paper but extends it significantly to include "
            "multi-modal architectures and instruction-tuned models. Our key finding is that the optimal "
            "compute allocation changes substantially when fine-tuning is included in the training budget."
        ),
        (
            "Motivation & Problem Statement",
            "Why do we care about scaling laws?",
            "The fundamental question we set out to answer is: given a fixed compute budget, what is the "
            "optimal way to allocate it between model size and training tokens? Previous work by Hoffmann "
            "et al. showed that most large models are significantly undertrained. Our work revisits this "
            "finding with a broader experimental scope, including models ranging from 70M to 70B parameters "
            "trained on datasets from 1B to 2T tokens. We find the compute-optimal frontier has shifted "
            "significantly with modern tokenizers and data curation pipelines."
        ),
        (
            "Related Work",
            "Standing on the shoulders of giants",
            "Key prior work includes Kaplan et al. 2020, which first established power-law scaling for "
            "language models. The Chinchilla paper by Hoffmann et al. 2022 was a major correction, showing "
            "models should be trained on roughly 20 tokens per parameter. The GPT-4 technical report "
            "provided empirical evidence but withheld training details. Our work differs in three ways: "
            "we use a more diverse model family, we include post-training in the compute budget, and we "
            "measure scaling on downstream task performance rather than only perplexity."
        ),
        (
            "Experimental Setup",
            "Models, datasets, and compute allocation",
            "We trained 412 models spanning seven orders of magnitude in compute, from 6×10^17 to "
            "6×10^24 FLOPs. All models use a transformer decoder architecture with RoPE positional "
            "embeddings and SwiGLU activations. Training data is a 2T token corpus drawn from "
            "Common Crawl (60%), GitHub (10%), Wikipedia (8%), Books (7%), ArXiv (5%), and other "
            "high-quality sources. We used a unified tokenizer with 128k vocabulary. Each experiment "
            "was run three times with different random seeds; we report median validation loss."
        ),
        (
            "Methodology: Measuring Scaling",
            "How we fit the scaling coefficients",
            "We parameterize loss as L(N, D) = E + A/N^alpha + B/D^beta, where N is model parameters, "
            "D is training tokens, E is irreducible entropy, and alpha, beta are fitted exponents. "
            "Fitting is performed via L-BFGS on the held-out validation set. A critical methodological "
            "improvement over prior work is our use of a held-out distribution shift test: we measure "
            "not just in-distribution perplexity but also zero-shot accuracy on 57 downstream tasks "
            "spanning reasoning, factual recall, coding, and mathematical problem solving."
        ),
        (
            "Results: Compute-Optimal Frontier",
            "N* = 0.31 × C^0.54, D* = 5.4 × C^0.46",
            "Our main result is the compute-optimal frontier: for a compute budget C FLOPs, the optimal "
            "model size is N* = 0.31 × C^0.54 parameters and the optimal token count is D* = 5.4 × C^0.46. "
            "This implies that at 10^23 FLOPs, the optimal model is approximately 67B parameters trained "
            "on 1.3T tokens — somewhat larger than the Chinchilla prediction. The difference arises because "
            "our data curation pipeline is more aggressive, and because we account for repeated data epochs. "
            "Importantly, this frontier shifts further when instruction tuning compute is included."
        ),
        (
            "Results: Downstream Task Scaling",
            "Benchmark performance vs. compute",
            "When measuring scaling on downstream tasks rather than perplexity, we observe a phase "
            "transition phenomenon: performance remains flat for small models then jumps sharply. "
            "This behavior is most pronounced on tasks requiring multi-step reasoning, such as "
            "MATH and GSM8K. The inflection point corresponds to approximately 7B parameters. "
            "Below this threshold, additional compute yields minimal downstream gains regardless "
            "of allocation strategy. Above it, our compute-optimal frontier consistently outperforms "
            "both overparameterized and undertrained baselines by 8-15% on our benchmark suite."
        ),
        (
            "Results: Data Quality vs. Quantity",
            "Quality filtering changes the optimal frontier",
            "We conducted an ablation study varying data quality at fixed compute. Using aggressive "
            "quality filtering (removing near-duplicate documents, low-perplexity outliers, and "
            "documents with less than 200 words) shifts the optimal frontier toward larger models "
            "per token. Intuitively, higher-quality data is more informative per token, so fewer "
            "tokens are needed. With our highest-quality 200B token subset, the compute-optimal "
            "model at 10^23 FLOPs grows to 180B parameters — a 2.7× increase over the unfiltered "
            "baseline. This has major implications for data curation as a scaling lever."
        ),
        (
            "Analysis: Why Perplexity Diverges from Downstream",
            "The perplexity trap",
            "A key finding of our work is that perplexity and downstream task accuracy can diverge "
            "substantially, particularly when data distribution shifts occur between training and "
            "evaluation. We term this the 'perplexity trap': a model optimized purely for perplexity "
            "on a fixed distribution may underperform on real-world tasks. We provide a theoretical "
            "analysis using information-theoretic arguments showing that perplexity is an upper bound "
            "on downstream task performance only when the evaluation distribution is contained within "
            "the training distribution. For out-of-distribution tasks, no such guarantee exists."
        ),
        (
            "Analysis: Instruction Tuning and Compute",
            "How fine-tuning changes the optimal allocation",
            "When we include instruction tuning (RLHF or SFT) in the total compute budget, the "
            "optimal pre-training strategy changes. Specifically, the optimal model for a given "
            "total budget (pre-training + fine-tuning) is smaller than the compute-optimal pre-trained "
            "model alone. This is because fine-tuning provides strong signal for aligning the model "
            "with user intent, partially compensating for reduced pre-training capacity. Our analysis "
            "suggests the ideal split is approximately 85% pre-training, 15% fine-tuning for "
            "instruction-following tasks. For specialized domains, fine-tuning share can increase to 30%."
        ),
        (
            "Limitations",
            "What our analysis does not capture",
            "Several important limitations must be acknowledged. First, our scaling laws are fit on "
            "a single model family and may not transfer to architectures with different inductive "
            "biases, such as mixture-of-experts or state-space models. Second, we focus on autoregressive "
            "pre-training; results for masked language models or diffusion-based text models may differ "
            "substantially. Third, our analysis treats all compute as equivalent, ignoring communication "
            "overhead in distributed training. In practice, pipeline parallelism can shift the "
            "effective compute budget by 15-20%. Finally, data diversity and quality effects remain "
            "underspecified in our current framework."
        ),
        (
            "Broader Impact",
            "What this means for the field",
            "Our findings have several implications for the broader research community. First, the "
            "shifted compute-optimal frontier means that many currently deployed models are either "
            "overtrained (too many tokens for their size) or under-sized. Second, our data quality "
            "findings suggest that the returns to data curation are much higher than previously "
            "estimated, which has implications for how organizations should allocate data engineering "
            "resources. Third, the instruction-tuning compute split has direct consequences for "
            "practitioners deciding how to allocate GPU budgets. We provide an open-source calculator "
            "to help teams apply our scaling laws to their specific settings."
        ),
        (
            "Future Work",
            "Open questions and next steps",
            "Several directions are promising for future investigation. First, extending our framework "
            "to multimodal models is a natural next step; preliminary results suggest scaling laws "
            "differ substantially for vision-language models. Second, the phase transition phenomenon "
            "we observe on reasoning tasks warrants deeper theoretical analysis; we conjecture it "
            "relates to the emergence of chain-of-thought reasoning capabilities. Third, scaling "
            "laws for reinforcement learning from human feedback are largely unexplored. Fourth, "
            "understanding how scaling laws interact with sparse architectures like MoE models "
            "could unlock significant efficiency gains. We are actively pursuing all four directions."
        ),
        (
            "Conclusion",
            "Key takeaways",
            "To summarize the main contributions of this paper. First, we establish an updated "
            "compute-optimal frontier that accounts for modern data curation and tokenization: "
            "N* = 0.31 × C^0.54, slightly larger than Chinchilla at the same compute budget. "
            "Second, we show that perplexity and downstream accuracy diverge when distribution "
            "shift is present, and recommend using downstream benchmarks for scaling law evaluation. "
            "Third, we show that including instruction tuning in the compute budget shifts the "
            "optimal pre-training allocation toward smaller, more heavily fine-tuned models. "
            "These findings should help practitioners make better-informed decisions about model "
            "training at scale."
        ),
        (
            "Acknowledgments & References",
            "Thank you",
            "We thank the reviewers for their insightful feedback. Computational resources were "
            "provided by the National Science Foundation under grant NSF-2134819 and by an industry "
            "partnership with three major cloud providers. Data engineering support was provided by "
            "the data team at our institution. We also thank colleagues in the NLP reading group "
            "for valuable discussions throughout this project. The full reference list is available "
            "in the paper. Code and trained model checkpoints will be released upon publication "
            "at our project page. Questions can be directed to the corresponding author."
        ),
    ]

    for idx, (title_text, content_text, notes_text) in enumerate(slides_data):
        if idx == 0:
            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            try:
                slide.placeholders[1].text = content_text
            except (KeyError, IndexError):
                pass
        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            try:
                slide.placeholders[1].text = content_text
            except (KeyError, IndexError):
                pass

        # Add presenter notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open NeurIPS_Talk.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
