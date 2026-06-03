"""
FINAL REWARD SCRIPT - SUCCESS
Task: Change the list marker of paragraphs 36–38 to square bullets.
Generated: 2025-10-17 16:13:09
Status: success
Model: azure-o3
Total Steps: 11
"""

import os
from pptx import Presentation
from pptx.shapes.group import GroupShape

"""
Reward Script: Verify that the list marker (bullet character) of paragraphs 36–38
in the provided PPTX file have been changed to SQUARE bullets (Unicode U+25AA),
while the bullets of paragraphs 1–35 remain the original ROUND bullets (Unicode
U+2022).

Scoring (progressive):
  • 0.7 points – paragraphs 36, 37, and 38 all use a square bullet
    (partial credit: 0.233 per correct paragraph).
  • 0.3 points – bullets for paragraphs 1-35 are unchanged from the original
    bullet used in paragraph 1 (partial credit proportional to how many remain
    unchanged).
The script prints detailed diagnostics and finally the reward as
```
REWARD: X.X
```
Exactly 1.0 is printed only when the task is perfectly completed.
"""

# ---------------------------------------------------------------------------
# Utility helpers
# ---------------------------------------------------------------------------

def _gather_all_shapes(container):
    """Recursively collect shapes from a slide / group shape."""
    shapes = []
    for shp in container.shapes:
        shapes.append(shp)
        if isinstance(shp, GroupShape):
            shapes.extend(_gather_all_shapes(shp))
    return shapes


def _gather_bullet_paragraphs(prs):
    """Return a list of tuples (paragraph, bullet_char) for every paragraph that
    explicitly has a bullet character assigned (buChar element present)."""
    bullet_paras = []
    for slide in prs.slides:
        for shape in _gather_all_shapes(slide):
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    bu_char_el = para._pPr.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
                    )
                    if bu_char_el is not None and bu_char_el.get("char") is not None:
                        bullet_paras.append((para, bu_char_el.get("char")))
    return bullet_paras


def _is_square(ch):
    """Return True if the supplied character represents a square bullet."""
    if ch is None:
        return False
    # Exact char
    if len(ch) == 1 and ord(ch) == 0x25AA:
        return True
    # Accept common variants / entities
    if ch in {"■", "▪", "\u25AA", "&#9642;"}:
        return True
    return False


def _is_same_char(ch, reference):
    """Safely compare bullet characters (handles entity vs literal)."""
    if ch == reference:
        return True
    # Compare by code-point when both are single characters
    if ch and reference and len(ch) == 1 and len(reference) == 1:
        return ord(ch) == ord(reference)
    return False

# ---------------------------------------------------------------------------
# Main verification function
# ---------------------------------------------------------------------------

def verify_bullet_task(file_path: str) -> float:
    MAX_SCORE = 1.0
    score = 0.0

    # --- Basic file checks --------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0
    if not file_path.lower().endswith(".pptx"):
        print("✗ Provided file is not a .pptx")
        return 0.0

    # --- Load presentation --------------------------------------------------
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to load presentation: {exc}")
        return 0.0

    # --- Collect bullet paragraphs -----------------------------------------
    bullet_paras = _gather_bullet_paragraphs(prs)
    total_bullets = len(bullet_paras)
    print(f"Total bullet paragraphs detected: {total_bullets}")

    if total_bullets < 38:
        print("✗ Fewer than 38 bullet paragraphs – cannot perform verification")
        return 0.0

    # --- Determine the original bullet character (para 1) -------------------
    original_bullet_char = bullet_paras[0][1]
    print(f"Original bullet character detected as: {repr(original_bullet_char)}")

    # --- Requirement 1: paragraphs 36-38 use square bullets -----------------
    target_indices = [36, 37, 38]  # 1-based indices in bullet_paras list
    correct_square = 0
    for idx in target_indices:
        para, ch = bullet_paras[idx - 1]
        print(f"Bullet #{idx}: text='{para.text}', bullet={repr(ch)}")
        if _is_square(ch):
            correct_square += 1
        else:
            print("  ✗ Expected square bullet here")

    square_score = (correct_square / 3) * 0.7  # up to 0.7 pts
    score += square_score
    if square_score == 0.7:
        print("✓ All target bullets correctly changed to square (0.7 points)")
    else:
        print(f"Partial square completion: {correct_square}/3 correct ({square_score:.2f} pts)")

    # --- Requirement 2: bullets 1-35 remain unchanged -----------------------
    unchanged = 0
    for idx in range(1, 36):
        para, ch = bullet_paras[idx - 1]
        if _is_same_char(ch, original_bullet_char):
            unchanged += 1
        else:
            print(f"  ✗ Bullet #{idx} altered (expected {repr(original_bullet_char)}, got {repr(ch)})")

    unchanged_score = (unchanged / 35) * 0.3  # up to 0.3 pts
    score += unchanged_score
    if unchanged_score == 0.3:
        print("✓ Bullets 1–35 remain unchanged (0.3 points)")
    else:
        print(f"Bullets unchanged: {unchanged}/35 ({unchanged_score:.2f} pts)")

    # --- Final score --------------------------------------------------------
    final_score = min(round(score, 4), MAX_SCORE)
    print(f"REWARD: {final_score}")
    return final_score

# ---------------------------------------------------------------------------
# Execute verification when script is run directly
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    file_path = "/home/user/change_the_list_marker_of_paragraphs_3638_to_square_bullets.pptx"
    verify_bullet_task(file_path)

