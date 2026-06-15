"""
Reward Script: Prepare deck for executive review
Task ID: osworld_impress_multi_op_combined_014
Domain: libreoffice_impress
Scoring:
  Component 1: Titles on slides 3, 4, 5 formatted bold + red + underline (0.5 pts total)
  Component 2: White background on all 8 slides (0.3 pts)
  Component 3: 'FOR INTERNAL USE ONLY' text removed from slide 1 (0.2 pts)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_014'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Task requirements:
    - Apply bold + red + underline to titles on slides 3, 4, and 5
    - Change the global background to white on all slides
    - Remove the 'FOR INTERNAL USE ONLY' text from slide 1
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have 8 slides
    if len(prs.slides) < 8:
        print(f"CRITICAL: Expected 8 slides, found {len(prs.slides)}. File may be corrupted.")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Title formatting on slides 3, 4, 5 (bold + red + underline)
    # Each slide's title contributes 0.5/3 ~ 0.167 points
    # We verify: bold=True, underline=True, and color=FF0000 (red)
    # Initial state: bold=False, underline=False, color=000000 — so this FAILS on initial
    slide_indices = [2, 3, 4]  # 0-based indices for slides 3, 4, 5
    slide_names = {2: 'slide 3', 3: 'slide 4', 4: 'slide 5'}
    title_texts = {2: 'Financial Performance', 3: 'Sales & Revenue Highlights', 4: 'Product Roadmap'}

    points_per_title = round(0.5 / 3, 4)  # ~0.1667 per slide

    for idx in slide_indices:
        try:
            slide = prs.slides[idx]
            # Find the title shape (first placeholder with title-like text)
            title_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip() in title_texts.values():
                    # Check if this is one of the expected titles
                    text = shape.text_frame.text.strip()
                    if text == title_texts[idx]:
                        title_shape = shape
                        break

            if title_shape is None:
                print(f"FAIL: Component 1 — Title shape not found on {slide_names[idx]}")
                continue

            # Get the runs with actual text from the title
            all_runs = []
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        all_runs.append(run)

            if not all_runs:
                print(f"FAIL: Component 1 — No text runs found in title on {slide_names[idx]}")
                continue

            # Check that ALL non-empty runs have bold=True, underline=True, color=FF0000
            title_ok = True
            fail_reasons = []

            for run in all_runs:
                # Check bold
                bold = run.font.bold
                if bold is not True:
                    title_ok = False
                    fail_reasons.append(f"bold={bold} (expected True)")

                # Check underline
                underline = run.font.underline
                if underline is not True:
                    title_ok = False
                    fail_reasons.append(f"underline={underline} (expected True)")

                # Check color is red (FF0000)
                try:
                    if run.font.color.type is not None:
                        rgb = str(run.font.color.rgb)
                        if rgb.upper() != 'FF0000':
                            title_ok = False
                            fail_reasons.append(f"color={rgb} (expected FF0000)")
                    else:
                        title_ok = False
                        fail_reasons.append("color=inherited (expected FF0000)")
                except Exception as ce:
                    title_ok = False
                    fail_reasons.append(f"color check error: {ce}")

            if title_ok:
                print(f"PASS: Component 1 — {slide_names[idx]} title '{title_texts[idx]}' is bold+red+underline ({points_per_title:.4f} pts)")
                total_score += points_per_title
            else:
                print(f"FAIL: Component 1 — {slide_names[idx]} title '{title_texts[idx]}': {'; '.join(fail_reasons)}")

        except Exception as e:
            print(f"ERROR: Component 1 — checking {slide_names[idx]}: {e}")

    # Component 2: White background on all 8 slides (0.3 points)
    # Initial state: all slides have background E8E8E8 (light gray)
    # Golden state: all slides have background FFFFFF (white)
    try:
        slides_with_white_bg = 0
        total_slides = len(prs.slides)
        fail_details = []

        for i, slide in enumerate(prs.slides):
            bg_fill = slide.background.fill
            if bg_fill.type == 1:  # SOLID fill
                try:
                    rgb = str(bg_fill.fore_color.rgb).upper()
                    if rgb == 'FFFFFF':
                        slides_with_white_bg += 1
                    else:
                        fail_details.append(f"slide {i+1}: color={rgb}")
                except Exception as ce:
                    fail_details.append(f"slide {i+1}: color read error {ce}")
            else:
                fail_details.append(f"slide {i+1}: fill type={bg_fill.type} (not solid)")

        if slides_with_white_bg == total_slides:
            print(f"PASS: Component 2 — All {total_slides} slides have white (FFFFFF) background (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — {slides_with_white_bg}/{total_slides} slides have white background; failures: {fail_details}")

    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: 'FOR INTERNAL USE ONLY' text removed from slide 1 (0.2 points)
    # Initial state: slide 1 has a text box with 'FOR INTERNAL USE ONLY'
    # Golden state: that text box is absent
    try:
        slide1 = prs.slides[0]
        found_internal_use_text = False
        for shape in slide1.shapes:
            if shape.has_text_frame:
                full_text = shape.text_frame.text.strip()
                if 'FOR INTERNAL USE ONLY' in full_text.upper():
                    found_internal_use_text = True
                    break

        if not found_internal_use_text:
            print(f"PASS: Component 3 — 'FOR INTERNAL USE ONLY' text is absent from slide 1 (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 3 — 'FOR INTERNAL USE ONLY' text still present on slide 1")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
