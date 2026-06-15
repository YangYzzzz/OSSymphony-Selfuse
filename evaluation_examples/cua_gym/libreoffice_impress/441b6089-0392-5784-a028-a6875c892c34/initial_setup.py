"""
Initial Setup: Executive Briefing Deck - Pre-Task State
Task ID: osworld_impress_multi_op_combined_014
Domain: libreoffice_impress

Creates an 8-slide executive briefing presentation with:
  - Light gray background on ALL slides
  - Plain black titles on slides 3-5 (no bold/red/underline)
  - 'FOR INTERNAL USE ONLY' text box on slide 1
  - Realistic business content throughout
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_014'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# Light gray background color
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
# Black text color for titles
BLACK = RGBColor(0x00, 0x00, 0x00)


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_slide_background(slide, color: RGBColor):
    """Set a solid background fill color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text: str, font_size: int = 32, bold: bool = False,
                   color: RGBColor = None, underline: bool = False):
    """Add styled text to the slide title placeholder."""
    title_shape = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            title_shape = shape
            break
    if title_shape is None:
        return
    title_shape.text = text
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.underline = underline
            if color:
                run.font.color.rgb = color
            else:
                run.font.color.rgb = BLACK


def add_content_text(slide, text: str, font_size: int = 18):
    """Add text to the content placeholder if available."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            return


def create_initial():
    prs = Presentation()
    # Use standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide with "FOR INTERNAL USE ONLY" ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background(slide1, LIGHT_GRAY)
    slide1.shapes.title.text = "Q1 2025 Executive Briefing"
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(40)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x6A)

    # Subtitle
    if len(slide1.placeholders) > 1:
        slide1.placeholders[1].text = "Strategic Overview — Acme Corporation\nMarch 2025"
        for para in slide1.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

    # Add "FOR INTERNAL USE ONLY" text box — MUST be present in initial state
    txBox = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(4.5), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FOR INTERNAL USE ONLY"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide2, LIGHT_GRAY)
    slide2.shapes.title.text = "Agenda"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x6A)

    if len(slide2.placeholders) > 1:
        agenda_text = (
            "1. Financial Performance\n"
            "2. Sales & Revenue\n"
            "3. Product Roadmap\n"
            "4. Customer Insights\n"
            "5. Operational Updates\n"
            "6. Q2 Outlook"
        )
        slide2.placeholders[1].text = agenda_text
        for para in slide2.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 3: Financial Performance — plain black title (NO bold/red/underline) ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide3, LIGHT_GRAY)
    slide3.shapes.title.text = "Financial Performance"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = False
            run.font.underline = False
            run.font.color.rgb = BLACK

    if len(slide3.placeholders) > 1:
        slide3.placeholders[1].text = (
            "Total Revenue: $48.7M (up 12% YoY)\n"
            "Gross Margin: 64.3%\n"
            "Operating Income: $9.2M\n"
            "EBITDA: $11.4M\n"
            "Cash & Equivalents: $22.1M\n"
            "Accounts Receivable: $7.8M"
        )
        for para in slide3.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 4: Sales & Revenue — plain black title (NO bold/red/underline) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide4, LIGHT_GRAY)
    slide4.shapes.title.text = "Sales & Revenue Highlights"
    for para in slide4.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = False
            run.font.underline = False
            run.font.color.rgb = BLACK

    if len(slide4.placeholders) > 1:
        slide4.placeholders[1].text = (
            "North America: $27.3M (+15%)\n"
            "EMEA: $13.6M (+8%)\n"
            "APAC: $7.8M (+22%)\n"
            "New Customers: 142\n"
            "Customer Retention Rate: 91.4%\n"
            "Average Deal Size: $342,500"
        )
        for para in slide4.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 5: Product Roadmap — plain black title (NO bold/red/underline) ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide5, LIGHT_GRAY)
    slide5.shapes.title.text = "Product Roadmap"
    for para in slide5.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = False
            run.font.underline = False
            run.font.color.rgb = BLACK

    if len(slide5.placeholders) > 1:
        slide5.placeholders[1].text = (
            "Q1 2025: Launch v3.2 with AI-assisted analytics\n"
            "Q2 2025: Mobile app redesign (iOS & Android)\n"
            "Q3 2025: Enterprise SSO and compliance module\n"
            "Q4 2025: API marketplace beta release\n"
            "Pending: Integration with Salesforce and HubSpot"
        )
        for para in slide5.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 6: Customer Insights ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide6, LIGHT_GRAY)
    slide6.shapes.title.text = "Customer Insights"
    for para in slide6.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x6A)

    if len(slide6.placeholders) > 1:
        slide6.placeholders[1].text = (
            "NPS Score: 68 (industry avg: 42)\n"
            "Support Tickets Resolved: 98.7% within SLA\n"
            "Churn Rate: 3.2% (down from 4.8%)\n"
            "Top Feedback: Better reporting tools (38%)\n"
            "Case Study: Meridian Health — 35% cost reduction"
        )
        for para in slide6.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 7: Operational Updates ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide7, LIGHT_GRAY)
    slide7.shapes.title.text = "Operational Updates"
    for para in slide7.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x6A)

    if len(slide7.placeholders) > 1:
        slide7.placeholders[1].text = (
            "Headcount: 387 employees (+ 42 since Q4)\n"
            "Engineering: Completed migration to Kubernetes\n"
            "HR: New benefits package effective April 1\n"
            "Facilities: Austin office expansion completed\n"
            "ISO 27001 recertification: In progress"
        )
        for para in slide7.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ---- Slide 8: Q2 Outlook ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide8, LIGHT_GRAY)
    slide8.shapes.title.text = "Q2 2025 Outlook"
    for para in slide8.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x6A)

    if len(slide8.placeholders) > 1:
        slide8.placeholders[1].text = (
            "Revenue Target: $52–55M\n"
            "Planned Hires: 28 (primarily Engineering & Sales)\n"
            "Priority: Close 3 enterprise deals in pipeline\n"
            "R&D Investment: Increase by 18% QoQ\n"
            "Board Meeting: April 28, 2025"
        )
        for para in slide8.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
