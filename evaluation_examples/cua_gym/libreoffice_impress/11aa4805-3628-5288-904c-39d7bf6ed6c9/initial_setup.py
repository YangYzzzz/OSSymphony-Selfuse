"""
Initial Setup: Create Implementation_Pitch.pptx with 10 slides.
Slide 8 has title 'Implementation Timeline' but empty content area.
Task ID: impress_sales_062
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_062'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text):
    """Set the title placeholder text."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_body_text(slide, lines):
    """Add lines of text to the first body placeholder (index 1)."""
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Enterprise CRM Implementation Pitch"
    slide1.placeholders[1].text = "Prepared for Meridian Healthcare Group\nQ2 2025"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "Executive Summary")
    add_body_text(slide2, [
        "Meridian Healthcare currently manages 12,000+ patient relationships across 3 regional offices",
        "Existing tools are fragmented: spreadsheets, email threads, and a legacy database from 2014",
        "A modern CRM platform will unify patient engagement tracking and reduce admin overhead by 35%",
        "Projected ROI: 280% within 18 months of full deployment",
    ])

    # --- Slide 3: Current Challenges ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide3, "Current Challenges")
    add_body_text(slide3, [
        "No single source of truth for patient contact history",
        "Average response time to patient inquiries: 48 hours",
        "Manual data entry consumes 15 hours/week per coordinator",
        "Compliance tracking relies on paper-based audit trails",
        "30% of referrals are lost due to follow-up gaps",
    ])

    # --- Slide 4: Proposed Solution ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide4, "Proposed Solution: Unified CRM Platform")
    add_body_text(slide4, [
        "Cloud-based CRM with HIPAA-compliant data handling",
        "Automated patient journey tracking from referral to follow-up",
        "Integrated appointment scheduling and reminder system",
        "Real-time dashboards for regional managers",
        "Mobile access for field coordinators",
    ])

    # --- Slide 5: Key Features ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide5, "Key Features & Capabilities")
    add_body_text(slide5, [
        "360-degree patient profiles with interaction history",
        "Automated email and SMS campaign workflows",
        "Custom reporting: referral pipeline, conversion rates, retention metrics",
        "Role-based access control for clinical and admin staff",
        "Integration with existing EHR system (Epic)",
    ])

    # --- Slide 6: Cost Analysis ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide6, "Cost Analysis")
    add_body_text(slide6, [
        "Software licensing: $4,200/month (50 seats)",
        "Implementation services: $38,000 one-time",
        "Data migration and cleanup: $12,500",
        "Staff training program: $8,000",
        "Total Year 1 investment: $108,900",
        "Projected annual savings: $185,000 (admin efficiency + referral capture)",
    ])

    # --- Slide 7: Success Metrics ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide7, "Success Metrics & KPIs")
    add_body_text(slide7, [
        "Patient response time reduced to < 4 hours",
        "Referral conversion rate increased from 45% to 70%",
        "Admin data entry time reduced by 60%",
        "100% compliance audit trail coverage",
        "Staff adoption rate > 90% within 60 days of launch",
    ])

    # --- Slide 8: Implementation Timeline (EMPTY CONTENT) ---
    # Use Title Only layout (index 5 = Blank, but we want a title)
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a textbox at the top
    txBox = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Implementation Timeline"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # NO other shapes - content area is intentionally empty

    # --- Slide 9: Risk Mitigation ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide9, "Risk Mitigation Strategy")
    add_body_text(slide9, [
        "Phased rollout minimizes operational disruption",
        "Dedicated project manager assigned for 8-week implementation",
        "Parallel running with legacy system during migration phase",
        "Weekly stakeholder check-ins and progress reporting",
        "Rollback plan documented for each migration phase",
    ])

    # --- Slide 10: Next Steps ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide10, "Next Steps & Call to Action")
    add_body_text(slide10, [
        "Schedule discovery workshop with IT and clinical leads (Week 1)",
        "Finalize vendor contract and SLA terms",
        "Kick off data audit and cleanup project",
        "Appoint internal CRM champions from each department",
        "Target go-live date: August 15, 2025",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
