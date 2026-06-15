"""
Reward Script: Create Gantt-style timeline on slide 8 with colored rectangle phases and week labels
Task ID: impress_sales_062
Domain: libreoffice_impress
Scoring:
  Component 1: Colored rectangle bars for 5 phases with correct colors (0.50)
  Component 2: Phase label text boxes (0.25)
  Component 3: Week number labels along the bottom (0.25)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_062'

# Expected phases: name -> hex color
EXPECTED_PHASES = {
    'Discovery': '4CAF50',
    'Setup': '2196F3',
    'Migration': 'FF9800',
    'Training': '9C27B0',
    'Go-Live': 'F44336',
}

# Expected week labels
EXPECTED_WEEKS = ['Wk 1', 'Wk 2', 'Wk 3', 'Wk 4', 'Wk 5', 'Wk 6', 'Wk 7', 'Wk 8']


def get_fill_color_hex(shape):
    """Get the solid fill color of a shape as uppercase hex string, or None."""
    try:
        fill = shape.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # Slide 8 (0-indexed)

    # Collect all shapes on slide 8
    rectangles = []  # (text, color_hex)
    text_boxes = []  # text content of non-rectangle text shapes

    for shape in slide.shapes:
        # Check for rectangle auto-shapes (the Gantt bars)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            color_hex = get_fill_color_hex(shape)
            shape_text = shape.text.strip() if hasattr(shape, 'text') else ''
            rectangles.append((shape_text, color_hex))
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            txt = shape.text.strip() if hasattr(shape, 'text') else ''
            if txt:
                text_boxes.append(txt)

    print(f"INFO: Found {len(rectangles)} rectangle shapes on slide 8")
    print(f"INFO: Found {len(text_boxes)} text boxes on slide 8")
    print(f"INFO: Rectangle details: {rectangles}")
    print(f"INFO: Text box contents: {text_boxes}")

    # Component 1: Colored rectangle bars for 5 phases (0.50 points, 0.10 each)
    try:
        matched_phases = 0
        for phase_name, expected_color in EXPECTED_PHASES.items():
            found = False
            for rect_text, rect_color in rectangles:
                # Check if rectangle has matching color (primary check)
                # and optionally matching text
                if rect_color and rect_color.upper() == expected_color.upper():
                    found = True
                    break
            if found:
                print(f"PASS: Component 1 — Phase '{phase_name}' rectangle with color #{expected_color} found (0.10 pts)")
                matched_phases += 1
                total_score += 0.10
            else:
                print(f"FAIL: Component 1 — Phase '{phase_name}' rectangle with color #{expected_color} not found")

        print(f"INFO: Component 1 subtotal: {matched_phases}/5 phases matched ({matched_phases * 0.10:.2f} pts)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Phase label text boxes (0.25 points, 0.05 each)
    try:
        matched_labels = 0
        for phase_name in EXPECTED_PHASES.keys():
            # Look for text boxes containing the phase name (not inside rectangles)
            # Phase labels are separate text boxes positioned to the left of bars
            label_found = False
            for txt in text_boxes:
                if phase_name.lower() in txt.lower():
                    label_found = True
                    break
            if label_found:
                print(f"PASS: Component 2 — Phase label '{phase_name}' found (0.05 pts)")
                matched_labels += 1
                total_score += 0.05
            else:
                print(f"FAIL: Component 2 — Phase label '{phase_name}' not found in text boxes")

        print(f"INFO: Component 2 subtotal: {matched_labels}/5 labels matched ({matched_labels * 0.05:.2f} pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Week number labels along the bottom (0.25 points)
    try:
        found_weeks = 0
        for wk_label in EXPECTED_WEEKS:
            wk_found = False
            for txt in text_boxes:
                # Match week labels flexibly: "Wk 1", "Wk1", "Week 1", "W1"
                if wk_label.lower() in txt.lower() or wk_label.replace(' ', '').lower() in txt.replace(' ', '').lower():
                    wk_found = True
                    break
            if wk_found:
                found_weeks += 1

        if found_weeks == 8:
            print(f"PASS: Component 3 — All 8 week labels found (0.25 pts)")
            total_score += 0.25
        elif found_weeks >= 6:
            partial = round(0.25 * (found_weeks / 8), 2)
            print(f"PARTIAL: Component 3 — {found_weeks}/8 week labels found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Only {found_weeks}/8 week labels found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
