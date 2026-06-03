"""
Reward Script: Create a new 6-slide presentation for an architecture portfolio with images.
Task ID: osworld_impress_new_presentation_images_007
Domain: libreoffice_impress

Scoring Rubric:
  Component 1: Presentation has exactly 6 slides (0.3 points)
  Component 2: Each slide has the correct image from ~/portfolio/ (0.7 points total,
               split evenly as ~0.1167 pts per slide, rounded)
               - Slide 1: building_1.jpg
               - Slide 2: building_2.jpg
               - Slide 3: building_3.jpg
               - Slide 4: interior_1.jpg
               - Slide 5: interior_2.jpg
               - Slide 6: site_plan.png
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_new_presentation_images_007'
PORTFOLIO_DIR = '/home/user/portfolio'

# Expected image filenames per slide (1-indexed)
EXPECTED_IMAGES = [
    'building_1.jpg',   # slide 1
    'building_2.jpg',   # slide 2
    'building_3.jpg',   # slide 3
    'interior_1.jpg',   # slide 4
    'interior_2.jpg',   # slide 5
    'site_plan.png',    # slide 6
]


def load_portfolio_blobs():
    """Load the binary content of each expected portfolio image."""
    blobs = {}
    for fname in EXPECTED_IMAGES:
        fpath = os.path.join(PORTFOLIO_DIR, fname)
        if os.path.exists(fpath):
            with open(fpath, 'rb') as f:
                blobs[fname] = f.read()
        else:
            blobs[fname] = None
    return blobs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: load the presentation file
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: load portfolio image blobs for comparison
    portfolio_blobs = load_portfolio_blobs()
    missing_sources = [k for k, v in portfolio_blobs.items() if v is None]
    if missing_sources:
        print(f"WARN: Portfolio source images not found: {missing_sources}")

    # Component 1: Presentation has exactly 6 slides (0.3 points)
    # This FAILS on initial (no pptx exists), PASSES on golden
    try:
        num_slides = len(prs.slides)
        if num_slides == 6:
            print(f"PASS: Component 1 — Presentation has exactly 6 slides (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Expected 6 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Each slide has the correct image (0.7 points total, ~0.1167 per slide)
    # We check that each slide contains at least one picture shape and its binary blob
    # matches the expected portfolio image. Scored progressively (per-slide).
    per_slide_score = round(0.7 / 6, 4)  # ~0.1167 per slide

    for slide_idx, expected_fname in enumerate(EXPECTED_IMAGES):
        slide_num = slide_idx + 1
        try:
            if slide_idx >= len(prs.slides):
                print(f"FAIL: Component 2 Slide {slide_num} — Slide does not exist")
                continue

            slide = prs.slides[slide_idx]
            expected_blob = portfolio_blobs.get(expected_fname)

            # Find all picture shapes on the slide
            picture_shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

            if not picture_shapes:
                print(f"FAIL: Component 2 Slide {slide_num} — No images found on slide (expected {expected_fname})")
                continue

            # Check if any picture on the slide matches the expected portfolio image
            if expected_blob is None:
                # If we can't load the source, we can't compare blobs; skip this check
                print(f"WARN: Component 2 Slide {slide_num} — Cannot compare: source image {expected_fname} not found")
                continue

            # Use any() to avoid direct True assignment
            blob_match_found = any(
                s.image.blob == expected_blob
                for s in picture_shapes
            )

            if blob_match_found:
                print(f"PASS: Component 2 Slide {slide_num} — Correct image '{expected_fname}' found ({per_slide_score} pts)")
                total_score += per_slide_score
            else:
                # Report what blobs we actually found
                found_sizes = [len(s.image.blob) for s in picture_shapes]
                expected_size = len(expected_blob) if expected_blob else 'unknown'
                print(f"FAIL: Component 2 Slide {slide_num} — Image blob mismatch. "
                      f"Expected '{expected_fname}' ({expected_size} bytes), "
                      f"found blobs of size {found_sizes}")
        except Exception as e:
            print(f"ERROR: Component 2 Slide {slide_num} — {e}")

    final_score = round(min(total_score, 1.0), 4)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
