"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m doing a final polish on a long training presentation in LibreOffice Impress, and I just spotted an outlier: slide 184 is still using the default body typeface. How can I change every content text box on that one slide to Liberation Sans Narrow at exactly 20 pt?
Generated: 2025-09-10 18:30:23
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

"""
Reward Script for LibreOffice Impress / PPTX Task
Task: Ensure that every BODY/CONTENT text box on slide 184 uses
       – Font  : Liberation Sans Narrow
       – Size  : exactly 20 pt

Scoring Logic (progressive – max 1.0):
  • 0.5 pts for correct font family on each run
  • 0.5 pts for correct font size on each run
Score for each dimension = (# runs that pass) / (total body-runs)
Total reward  = 0.5*font_score + 0.5*size_score
Returns exactly 1.0 only if EVERY body-text run matches both criteria.

Anti-bias: No points for file existence, loading, or natural conditions.
"""

FILE_PATH = "/home/user/im_doing_a_final_polish_on_a_long_training_presentation_in_libreoffice_impress_and_i_just_spotted_an_golden.pptx"
SLIDE_INDEX = 183  # zero-based → slide 184 to user
EXPECTED_FONT = "Liberation Sans Narrow"
EXPECTED_SIZE_PT = 20

# ───────────────────────────────────────────────────────────────────────────────
# Helper functions
# ───────────────────────────────────────────────────────────────────────────────

def _get_run_font_name(run):
    """Return the run's font name, falling back to underlying XML if needed."""
    name = run.font.name
    if name:
        return name
    rPr = run._r.get_or_add_rPr()
    for tag in (qn("a:latin"), qn("a:ea"), qn("a:cs")):
        node = rPr.find(tag)
        if node is not None and node.get("typeface"):
            return node.get("typeface")
    # Last-ditch: attribute directly on rPr (rare in PPTX)
    return rPr.get(qn("a:typeface"))

def _get_run_font_size_pt(run):
    """Return the run's size in points, using XML fallback when necessary."""
    size = run.font.size
    if size is not None:
        return size.pt
    rPr = run._r.get_or_add_rPr()
    sz_attr = rPr.get("sz")  # stored in 1/100 pt
    if sz_attr and sz_attr.isdigit():
        return int(sz_attr) / 100.0
    return None

# ───────────────────────────────────────────────────────────────────────────────
# Core verification routine
# ───────────────────────────────────────────────────────────────────────────────

def verify_slide_body_font(file_path=FILE_PATH,
                            slide_index=SLIDE_INDEX,
                            expected_font=EXPECTED_FONT,
                            expected_size_pt=EXPECTED_SIZE_PT):
    print(f"Verifying slide {slide_index+1} in '{file_path}' …")

    # 0️⃣  Preliminary checks (no reward for passing these!)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to load presentation: {exc}")
        return 0.0
    if slide_index >= len(prs.slides):
        print(f"✗ Presentation has only {len(prs.slides)} slides – index {slide_index} invalid")
        return 0.0

    slide = prs.slides[slide_index]

    # 1️⃣  Collect all text runs from *body/content* placeholders (skip titles)
    total_runs = 0
    font_pass   = 0
    size_pass   = 0

    for shape in slide.shapes:
        # interested only in shapes with text
        if not shape.has_text_frame:
            continue
        # Skip TITLE placeholders explicitly (keep body/content boxes)
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue  # ignore empty runs
                total_runs += 1
                name = _get_run_font_name(run)
                size = _get_run_font_size_pt(run)
                print(f"Run: '{run.text.strip()}' | font='{name}' | size={size}")
                if name and name.lower() == expected_font.lower():
                    font_pass += 1
                if size is not None and abs(size - expected_size_pt) < 0.5:  # ±0.5 pt tolerance
                    size_pass += 1

    if total_runs == 0:
        print("✗ No body text found on target slide")
        return 0.0

    # 2️⃣  Progressive scoring
    font_score = font_pass / total_runs
    size_score = size_pass / total_runs
    reward = 0.5 * font_score + 0.5 * size_score

    # 3️⃣  Detailed report
    print(f"Font-family correct on {font_pass}/{total_runs} runs → {font_score:.2%}")
    print(f"Font-size   correct on {size_pass}/{total_runs} runs → {size_score:.2%}")
    print(f"TOTAL SCORE: {reward:.2f}")
    return reward

# ───────────────────────────────────────────────────────────────────────────────
# Execute verification when run as script
# ───────────────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    final_reward = verify_slide_body_font()
    print(f"REWARD: {final_reward}")
