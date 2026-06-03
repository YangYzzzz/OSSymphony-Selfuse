"""
Initial Setup: Keynote-style presentation with elaborate presenter notes
Task ID: osworld_multi_apps_impress_notes_export_014
Domain: libreoffice_impress (multi-app: also uses LibreOffice Writer)
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_notes_export_014'
DESKTOP = f'{WORKDIR}/Desktop'
OUTPUT_PPTX = f'{DESKTOP}/Keynote_Final.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Keynote slide data: (title, subtitle_or_content, presenter_notes)
SLIDES_DATA = [
    (
        "The Future of Artificial Intelligence",
        "A Visionary Keynote Address\nDr. Eleanor Vance, Chief Innovation Officer",
        "Welcome everyone. Take a deep breath and let the energy in the room sink in. "
        "Start by acknowledging the diversity of backgrounds represented today — engineers, "
        "ethicists, policymakers, and entrepreneurs. Pause briefly for effect before moving "
        "to slide two. Smile warmly and say: 'What we discuss today may well define the next "
        "century of human progress.'"
    ),
    (
        "Where We Stand Today",
        "AI Milestones of the Last Decade",
        "Reference the 2012 ImageNet breakthrough as the inflection point. Mention AlphaGo's "
        "victory in 2016 as an example of surpassing human intuition. Note GPT-4's multimodal "
        "capabilities. Use the phrase 'exponential is not a metaphor here — it is the data.' "
        "Pause for the audience to absorb the timeline graphic. Speak slowly and clearly."
    ),
    (
        "The Three Pillars of Modern AI",
        "Data  •  Compute  •  Algorithms",
        "Walk through each pillar individually. For Data: emphasize curation over volume — "
        "garbage in, garbage out. For Compute: acknowledge the carbon footprint conversation "
        "and how neuromorphic chips are a promising alternative. For Algorithms: highlight "
        "attention mechanisms and their role in transformers. Spend about 90 seconds per pillar."
    ),
    (
        "Healthcare Revolution",
        "From Diagnosis to Drug Discovery",
        "Open with the stat: AI-assisted radiology reduced missed diagnoses by 34% in a 2024 "
        "Stanford study. Transition to drug discovery — mention how AlphaFold solved the "
        "50-year protein folding problem. Share the anecdote about DeepMind's malaria research. "
        "Ask the audience: 'How many of you have a family member who could benefit from "
        "AI-accelerated medicine?' Let hands go up. Acknowledge them personally."
    ),
    (
        "Climate and Sustainability",
        "AI as a Tool for Planetary Health",
        "Begin with urgency — reference the IPCC 2023 report. Pivot to optimism: AI models "
        "predicting wildfire spread with 89% accuracy, energy grid optimization reducing "
        "waste by 15-20%. Mention Project Sunroof's ML-based solar mapping. Emphasize that "
        "AI is not a silver bullet but an accelerant for human-led solutions. End with: "
        "'Our planet gave us the conditions for intelligence — now intelligence must repay the debt.'"
    ),
    (
        "Education Reimagined",
        "Personalized Learning at Global Scale",
        "Describe the traditional one-size-fits-all classroom model and its limitations. "
        "Introduce adaptive learning platforms like Khan Academy's Khanmigo. Share the "
        "outcome data from Duolingo's AI-driven curriculum — 34% faster language acquisition. "
        "Address the equity concern head-on: 'What about students without reliable internet?' "
        "Mention offline AI initiatives in Sub-Saharan Africa. Conclude with the vision of "
        "every child having a personalized learning companion."
    ),
    (
        "The Workforce Transformation",
        "Disruption, Displacement, and Opportunity",
        "Acknowledge the anxiety in the room directly. Say: 'Yes, jobs will change. Some "
        "will disappear. But history teaches us that technology creates more than it destroys.' "
        "Reference the McKinsey 2024 Future of Work report — 12 million occupational transitions "
        "by 2030, offset by 20 million new roles. Emphasize reskilling. Share the story of "
        "a Toledo auto worker who became a robotics technician. Keep tone optimistic but honest."
    ),
    (
        "Ethical AI Frameworks",
        "Principles for Responsible Development",
        "This is a critical slide — slow down. List the core principles: Fairness, "
        "Accountability, Transparency, Safety (FATS framework). Discuss the EU AI Act as "
        "a landmark regulatory milestone. Reference the Montreal Declaration on Responsible AI. "
        "Ask: 'Who in this room has thought about the ethical implications of their AI product?' "
        "Pause. 'That number needs to be 100 percent.' Wait for it to land."
    ),
    (
        "Bias and Fairness",
        "The Hidden Variables in Machine Learning",
        "Open with the COMPAS recidivism algorithm controversy — a powerful cautionary tale. "
        "Explain how training data encodes historical inequalities. Introduce bias mitigation "
        "techniques: adversarial debiasing, reweighting, counterfactual fairness. Stress that "
        "fairness is not just a technical problem — it is a social and political one. "
        "Quote Timnit Gebru: 'You cannot technologize your way out of a social problem.'"
    ),
    (
        "Privacy and Surveillance",
        "The Double-Edged Blade of Recognition Technologies",
        "Begin with the statistic: over 1 billion surveillance cameras globally by 2025. "
        "Discuss facial recognition bans in San Francisco and Portland. Contrast with "
        "legitimate uses: missing persons, border security. Introduce differential privacy "
        "as a technical safeguard. Reference Apple's on-device processing model. Close with: "
        "'The question is not whether AI can watch us — it is whether we should let it, and who decides.'"
    ),
    (
        "The Geopolitics of AI",
        "National Strategies and Global Competition",
        "Frame the AI race as a strategic priority — US, China, EU, India each pursuing "
        "distinct national AI strategies. Reference China's 2030 AI supremacy plan. Discuss "
        "the semiconductor supply chain tension and TSMC's critical role. Introduce the "
        "concept of AI sovereignty. Urge for international cooperation frameworks, "
        "similar to nuclear non-proliferation treaties. Tone: measured and serious."
    ),
    (
        "Open Source vs. Closed Models",
        "The Democratization Debate",
        "Frame the tension: open-source accelerates innovation but lowers misuse barriers. "
        "Closed models offer safety controls but concentrate power. Reference Meta's LLaMA "
        "release and its effects. Contrast with OpenAI's graduated access model. "
        "Ask: 'Should the blueprint for the most powerful technology in history be public?' "
        "Allow a beat of silence. Present both sides fairly. Acknowledge this is an "
        "unresolved debate you are personally wrestling with."
    ),
    (
        "Generative AI and Creativity",
        "Partner, Tool, or Threat to Human Expression?",
        "Open with a live example if possible — show AI-generated music or art on screen. "
        "Discuss the copyright controversy: the New York Times v. OpenAI case. "
        "Introduce the concept of creative augmentation vs. replacement. Quote David Bowie: "
        "'The Internet is an alien life form.' Argue AI is a similar paradigm shift for art. "
        "End with: 'The canvas has changed. The artist has not.'"
    ),
    (
        "AI Safety and Existential Risk",
        "Long-Term Alignment Research",
        "Transition carefully — this is where tone becomes most serious. "
        "Reference Nick Bostrom's Superintelligence. Introduce the alignment problem: "
        "how do you ensure AI systems pursue human values as they become more capable? "
        "Discuss RLHF (Reinforcement Learning from Human Feedback). Mention Anthropic, "
        "DeepMind Safety, and OpenAI's superalignment team. Be honest: 'We do not have "
        "this solved. But the best minds in the world are working on it.' Steady voice."
    ),
    (
        "The Path to AGI",
        "Timelines, Milestones, and Uncertainty",
        "Survey expert predictions — from 5 years to never. Introduce the concept of "
        "narrow vs. general vs. superintelligent AI. Discuss emergent behaviors in large "
        "language models as a potential early signal. Reference the Turing Test and its "
        "limitations. Introduce more rigorous benchmarks: ARC-AGI, FrontierMath. "
        "Conclude: 'Whether AGI is 5 years or 50 years away, the ethical frameworks "
        "we build today will determine what it becomes.'"
    ),
    (
        "AI in Scientific Discovery",
        "Accelerating the Pace of Human Knowledge",
        "Open with the Nobel Prize acknowledgment of AlphaFold in 2024. Discuss AI's role "
        "in fusion energy research at Lawrence Livermore. Mention materials science breakthroughs "
        "from Graph Neural Networks. Describe GNoME's discovery of 2.2 million new crystal "
        "structures. Say: 'For the first time in history, we have a collaborator that never "
        "sleeps, never tires, and can read every paper ever written.' Pause. Let that vision "
        "resonate with the scientists in the room."
    ),
    (
        "The Human-AI Partnership",
        "Augmentation Over Replacement",
        "Shift to an empowering frame. Introduce centaur chess as a metaphor — human plus AI "
        "beats both humans alone and AI alone. Share examples from surgery, legal research, "
        "and financial analysis. Emphasize that the goal is to amplify human judgment, "
        "not replace it. Quote Garry Kasparov: 'Advanced chess showed us that weak human + "
        "machine + better process beats strong computer alone.' Build to: 'This is the future "
        "we should be designing for — not a race between humans and machines, but a collaboration.'"
    ),
    (
        "Policy Recommendations",
        "Building the Infrastructure for Beneficial AI",
        "Present five actionable recommendations: 1) Establish a Global AI Safety Institute. "
        "2) Require algorithmic impact assessments for high-stakes systems. 3) Fund AI literacy "
        "education from K-12 through continuing education. 4) Create adaptive regulatory "
        "sandboxes for AI experimentation. 5) Develop international AI governance treaties. "
        "Speak with authority and conviction here. You want this section to be quotable. "
        "Make eye contact with anyone in the audience who appears to be a policymaker."
    ),
    (
        "Call to Action",
        "Your Role in Shaping the AI Future",
        "Direct, personal tone. 'Every person in this room has a role to play.' Break it down "
        "by audience segment: engineers must build with ethics baked in; executives must "
        "demand responsible deployment; policymakers must lead without waiting for perfect "
        "information; educators must prepare students for jobs that do not yet exist; "
        "citizens must stay informed and engaged. Build emotional momentum. Voice should "
        "rise slightly. 'The future is not something that happens to us — it is something we build.'"
    ),
    (
        "Thank You",
        "Questions and Discussion",
        "Pause after stepping back from the podium. Take a breath. Say: 'I am deeply grateful "
        "for your time and attention today.' Acknowledge any co-speakers, sponsors, or "
        "organizing committee by name. Open Q&A with: 'I have covered a lot of ground today, "
        "and I have deliberately left some of the hardest questions unanswered — because I "
        "believe we need to answer them together. So — what questions do you have?' "
        "Smile. Be present. Listen fully before responding."
    ),
]


def create_initial():
    # Ensure Desktop directory exists
    os.makedirs(DESKTOP, exist_ok=True)

    prs = Presentation()
    # Standard widescreen dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, (title_text, content_text, notes_text) in enumerate(SLIDES_DATA):
        if i == 0:
            # Title slide layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            # Subtitle
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content_text
        else:
            # Title + Content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content_text

        # Add presenter notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = notes_text

    prs.save(OUTPUT_PPTX)
    print(f'Initial PPTX created: {OUTPUT_PPTX}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT_PPTX}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
