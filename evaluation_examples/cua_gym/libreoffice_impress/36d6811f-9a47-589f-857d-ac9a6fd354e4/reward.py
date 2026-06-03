"""
Reward Script: Extract presenter notes from Keynote_Final.pptx and create keynote_script.docx
Task ID: osworld_multi_apps_impress_notes_export_014
Domain: libreoffice_impress / libreoffice_writer (multi-app)
Scoring:
  Component 1 (0.25): keynote_script.docx has exactly 40 paragraphs (20 slides x 2 per slide)
  Component 2 (0.35): All 20 slide header paragraphs are centered with '--- Slide N ---' format
  Component 3 (0.40): Notes paragraphs contain actual slide notes content from Keynote_Final.pptx
  Total: 1.0
"""

import os
import re

# Domain-specific imports
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from pptx import Presentation

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_014'

DOCX_PATH = f'{WORKDIR}/keynote_script.docx'
PPTX_PATH = f'{WORKDIR}/Keynote_Final.pptx'
EXPECTED_SLIDE_COUNT = 20


def persist_app_state():
    """Send Ctrl+S to save any open LibreOffice documents."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for LibreOffice Writer")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task():
    """
    Verify that keynote_script.docx was created on the Desktop with:
    - 40 paragraphs total (2 per slide: header + notes body)
    - Each even paragraph is centered with '--- Slide N ---' format
    - Each odd paragraph contains the actual notes text from Keynote_Final.pptx
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: keynote_script.docx must exist
    if not os.path.exists(DOCX_PATH):
        print(f"FAIL: keynote_script.docx not found at {DOCX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Load the docx
    try:
        doc = Document(DOCX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load {DOCX_PATH}: {e}")
        print("REWARD: 0.0")
        return 0.0

    paras = doc.paragraphs
    actual_para_count = len(paras)

    # Precondition gate: Keynote_Final.pptx must exist to verify notes content
    if not os.path.exists(PPTX_PATH):
        print(f"WARN: Keynote_Final.pptx not found at {PPTX_PATH}. Notes content verification skipped.")
        prs = None
    else:
        try:
            prs = Presentation(PPTX_PATH)
        except Exception as e:
            print(f"WARN: Cannot load {PPTX_PATH}: {e}")
            prs = None

    # Extract slide notes from the pptx for comparison
    slide_notes = []
    if prs is not None:
        for slide in prs.slides:
            try:
                note = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                note = ''
            slide_notes.append(note)

    # -------------------------------------------------------------------------
    # Component 1: Document has exactly 40 paragraphs (20 slides x 2 per slide)
    # This FAILS on initial_env (no docx file at all → returns 0 early),
    # and PASSES on golden_env (40 paragraphs present).
    # -------------------------------------------------------------------------
    # Component 1: Correct paragraph count (0.25 points)
    try:
        expected_para_count = EXPECTED_SLIDE_COUNT * 2  # 40
        if actual_para_count == expected_para_count:
            print(f"PASS: Component 1 — Document has exactly {actual_para_count} paragraphs ({EXPECTED_SLIDE_COUNT} slides x 2) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected {expected_para_count} paragraphs, found {actual_para_count}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -------------------------------------------------------------------------
    # Component 2: All 20 slide header paragraphs are centered with '--- Slide N ---' format
    # Even-indexed paragraphs (0, 2, 4, ..., 38) should be centered headers.
    # -------------------------------------------------------------------------
    # Component 2: Centered slide headers with correct format (0.35 points)
    try:
        headers_correct = 0
        headers_found = 0

        for slide_idx in range(EXPECTED_SLIDE_COUNT):
            para_idx = slide_idx * 2
            if para_idx >= actual_para_count:
                break

            header_para = paras[para_idx]
            expected_header = f'--- Slide {slide_idx + 1} ---'
            actual_text = header_para.text.strip()
            actual_alignment = header_para.paragraph_format.alignment

            is_correct_text = (actual_text == expected_header)
            # CENTER alignment = WD_PARAGRAPH_ALIGNMENT.CENTER = 1
            is_centered = (actual_alignment == WD_PARAGRAPH_ALIGNMENT.CENTER)

            if is_correct_text and is_centered:
                headers_correct += 1
            else:
                if not is_correct_text:
                    print(f"FAIL: Component 2 — Slide {slide_idx + 1} header text: expected {repr(expected_header)}, got {repr(actual_text)}")
                if not is_centered:
                    print(f"FAIL: Component 2 — Slide {slide_idx + 1} header not centered (alignment={actual_alignment})")
            headers_found += 1

        if headers_correct == EXPECTED_SLIDE_COUNT:
            print(f"PASS: Component 2 — All {headers_correct}/{EXPECTED_SLIDE_COUNT} slide headers are correctly formatted and centered (0.35 pts)")
            total_score += 0.35
        else:
            partial = 0.35 * (headers_correct / EXPECTED_SLIDE_COUNT) if headers_correct >= EXPECTED_SLIDE_COUNT * 0.75 else 0.0
            if partial > 0.0:
                print(f"PARTIAL: Component 2 — {headers_correct}/{EXPECTED_SLIDE_COUNT} headers correct — awarding partial credit ({partial:.2f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 — Only {headers_correct}/{EXPECTED_SLIDE_COUNT} slide headers correct (need ≥75% for partial credit)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -------------------------------------------------------------------------
    # Component 3: Notes paragraphs contain actual slide notes content
    # Odd-indexed paragraphs (1, 3, 5, ..., 39) should contain notes from pptx.
    # We verify each notes paragraph is non-empty and its first 50 chars match
    # the actual slide notes from Keynote_Final.pptx.
    # -------------------------------------------------------------------------
    # Component 3: Notes content matches pptx slide notes (0.40 points)
    try:
        notes_correct = 0
        notes_checked = 0

        if prs is not None and len(slide_notes) == EXPECTED_SLIDE_COUNT:
            for slide_idx in range(EXPECTED_SLIDE_COUNT):
                notes_para_idx = slide_idx * 2 + 1
                if notes_para_idx >= actual_para_count:
                    break

                notes_para = paras[notes_para_idx]
                actual_notes = notes_para.text.strip()
                expected_notes = slide_notes[slide_idx].strip()

                # Check: non-empty and first 80 chars match the source slide notes
                is_nonempty = len(actual_notes) > 0
                # Use first 80 chars for matching (tolerates minor trailing differences)
                prefix_len = min(80, len(expected_notes))
                is_content_match = (
                    is_nonempty and
                    len(actual_notes) >= prefix_len and
                    actual_notes[:prefix_len] == expected_notes[:prefix_len]
                )

                if is_content_match:
                    notes_correct += 1
                else:
                    if not is_nonempty:
                        print(f"FAIL: Component 3 — Slide {slide_idx + 1} notes paragraph is empty")
                    else:
                        print(f"FAIL: Component 3 — Slide {slide_idx + 1} notes content mismatch")
                        print(f"  expected start: {repr(expected_notes[:60])}")
                        print(f"  actual start:   {repr(actual_notes[:60])}")
                notes_checked += 1

            if notes_correct == EXPECTED_SLIDE_COUNT:
                print(f"PASS: Component 3 — All {notes_correct}/{EXPECTED_SLIDE_COUNT} slide notes paragraphs match pptx source content (0.40 pts)")
                total_score += 0.40
            else:
                partial = 0.40 * (notes_correct / EXPECTED_SLIDE_COUNT) if notes_correct >= EXPECTED_SLIDE_COUNT * 0.75 else 0.0
                if partial > 0.0:
                    print(f"PARTIAL: Component 3 — {notes_correct}/{EXPECTED_SLIDE_COUNT} notes correct — awarding partial ({partial:.2f} pts)")
                    total_score += partial
                else:
                    print(f"FAIL: Component 3 — Only {notes_correct}/{EXPECTED_SLIDE_COUNT} notes paragraphs match (need ≥75% for partial credit)")
        else:
            # Fall back: check that all notes paragraphs are non-empty
            nonempty_count = 0
            for slide_idx in range(EXPECTED_SLIDE_COUNT):
                notes_para_idx = slide_idx * 2 + 1
                if notes_para_idx >= actual_para_count:
                    break
                notes_para = paras[notes_para_idx]
                if notes_para.text.strip():
                    nonempty_count += 1

            if nonempty_count == EXPECTED_SLIDE_COUNT:
                print(f"PASS (fallback): Component 3 — All {nonempty_count}/{EXPECTED_SLIDE_COUNT} notes paragraphs are non-empty (0.40 pts)")
                total_score += 0.40
            else:
                print(f"FAIL (fallback): Component 3 — Only {nonempty_count}/{EXPECTED_SLIDE_COUNT} notes paragraphs are non-empty")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score:.1f}")
    return final_score


# Entry point: persist any open GUI state, then verify
persist_app_state()
verify_task()
