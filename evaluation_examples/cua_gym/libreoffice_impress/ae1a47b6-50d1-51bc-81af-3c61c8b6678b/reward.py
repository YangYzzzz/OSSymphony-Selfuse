"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert comma-separated values in paragraph 3 into a table (comma delimiter).
Generated: 2025-10-17 07:00:59
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

"""
Reward Script: Verify conversion of comma-separated values in paragraph 3 into a table
----------------------------------------------------------------------
Task to verify
  • Paragraph 3 originally contained a CSV string (comma-delimited)
  • The user had to convert that CSV paragraph into a table
Verification Logic
  1.   Ensure the CSV paragraph no longer appears (no paragraphs with ≥2 commas)
  2.   Ensure at least one table with more than one column now exists
  3.   Ensure none of the table cells still contain commas (data properly split)
Progressive Scoring (adds up to 1.0)
  • 0.3 – CSV paragraph removed
  • 0.4 – A multi-column table exists (shows the conversion took place)
  • 0.3 – All table cells free of commas (data truly split)
The script prints detailed diagnostics and finally prints
    REWARD: <score>
Only when every check passes does the score reach 1.0.
"""

def verify_csv_conversion_to_table(file_path: str) -> float:
    max_score = 1.0
    score = 0.0

    # --- 1. Load the presentation ---
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0  # Cannot continue without the file

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation; slide count: {len(prs.slides)}")
    except Exception as e:
        print(f"✗ Error loading PPTX: {e}")
        print("REWARD: 0.0")
        return 0.0

    # --- 2. Check for leftover CSV paragraph (>=2 commas) ---
    csv_paragraph_still_present = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                continue  # We only care about text paragraphs here
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text.count(",") >= 2:  # Likely old CSV text
                        csv_paragraph_still_present = True
                        print(f"✗ CSV-like paragraph still present: '{text[:60]}…'")
    if not csv_paragraph_still_present:
        score += 0.3
        print("✓ CSV paragraph removed (0.3 points)")

    # --- 3. Look for a multi-column table ---
    table_found = False
    table_cells_clean = True  # Assume clean until proven otherwise

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            cols = len(tbl.columns)
            rows = len(tbl.rows)
            if cols > 1:  # Conversion should yield >1 column
                table_found = True
                print(f"✓ Table found with {rows} rows x {cols} columns")
                # Check each cell for commas
                for r in range(rows):
                    for c in range(cols):
                        cell_text = tbl.cell(r, c).text_frame.text
                        if "," in cell_text:
                            table_cells_clean = False
                            print(f"   ✗ Comma found in table cell ({r},{c}): '{cell_text}'")
    if table_found:
        score += 0.4
        print("✓ Multi-column table present (0.4 points)")
    else:
        print("✗ No suitable table found (0 points)")

    if table_found and table_cells_clean:
        score += 0.3
        print("✓ Table cells free of commas (0.3 points)")
    elif table_found:
        print("✗ Commas still present in table cells (0 points)")

    final_score = min(score, max_score)
    print(f"TOTAL SCORE: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------------------------------------------------------------
# Execute verification when run as script
if __name__ == "__main__":
    FILE_PATH = "/home/user/convert_comma_separated_values_in_paragraph_3_into_a_table_comma_delimiter.pptx"
    verify_csv_conversion_to_table(FILE_PATH)

