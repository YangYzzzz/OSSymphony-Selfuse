"""
Reward Script: Apply striped table style to table on slide 5
Task ID: impress_gf1_043
Domain: libreoffice_impress
Scoring (only task-introduced changes):
  Component 1: Header row background teal + bold text — 0.30 pts
  Component 2: Header row white text WITH teal bg (compound) — 0.15 pts
  Component 3: Even data rows (2,4,6) have light teal #E0F2F1 — 0.30 pts
  Component 4: Outer borders removed (noFill) — 0.25 pts
Note: White row backgrounds and white font color exist in initial state,
so they are NOT scored independently. Only task-introduced changes score.
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_043'


def persist_app_state(domain: str):
    """Best-effort save for LibreOffice GUI."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.oxml.ns import qn
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)

    # Find the table on slide 5
    table = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    if table is None:
        print("CRITICAL: No table found on slide 5")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: table should be 7 rows x 5 cols
    if len(table.rows) != 7 or len(table.columns) != 5:
        print(f"CRITICAL: Table is {len(table.rows)}x{len(table.columns)}, expected 7x5")
        print("REWARD: 0.0")
        return 0.0

    num_rows = len(table.rows)
    num_cols = len(table.columns)

    def get_cell_fill_color(cell):
        """Get the srgbClr fill value of a cell, or None."""
        tcPr = cell._tc.find(qn('a:tcPr'))
        if tcPr is None:
            return None
        sf = tcPr.find(qn('a:solidFill'))
        if sf is None:
            return None
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            return srgb.get('val', '').upper()
        return None

    def get_border_info(cell, border_name):
        """Get border info: 'noFill', 'solid=COLOR', or 'absent'."""
        tcPr = cell._tc.find(qn('a:tcPr'))
        if tcPr is None:
            return 'absent'
        ln = tcPr.find(qn(f'a:{border_name}'))
        if ln is None:
            return 'absent'
        nf = ln.find(qn('a:noFill'))
        if nf is not None:
            return 'noFill'
        sf = ln.find(qn('a:solidFill'))
        if sf is not None:
            clr = sf.find(qn('a:srgbClr'))
            if clr is not None:
                return f"solid={clr.get('val', '').upper()}"
            return "solid=unknown"
        return 'other'

    # ===== Component 1: Header row has teal background AND bold text (0.30 pts) =====
    # Both teal bg and bold are task-introduced changes (initial has white bg, no bold)
    try:
        header_teal_bold_pass = 0
        for c in range(num_cols):
            cell = table.cell(0, c)
            fill = get_cell_fill_color(cell)
            is_teal = (fill == '00897B')

            is_bold = False
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        if run.font.bold is True:
                            is_bold = True
                        break
                break

            if is_teal and is_bold:
                header_teal_bold_pass += 1

        if header_teal_bold_pass == num_cols:
            print(f"PASS: Component 1 — All {num_cols} header cells have teal bg + bold text (0.30 pts)")
            total_score += 0.30
        elif header_teal_bold_pass > 0:
            partial = 0.30 * (header_teal_bold_pass / num_cols)
            print(f"PARTIAL: Component 1 — {header_teal_bold_pass}/{num_cols} header cells have teal bg + bold ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No header cells have both teal bg and bold text")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ===== Component 2: Header row white text WITH teal background (0.15 pts) =====
    # White text alone is pre-existing; we require teal bg as anchor to ensure this is task-introduced
    try:
        header_compound_pass = 0
        for c in range(num_cols):
            cell = table.cell(0, c)
            fill = get_cell_fill_color(cell)
            is_teal = (fill == '00897B')

            has_white_text = False
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        try:
                            if run.font.color.type is not None and str(run.font.color.rgb).upper() == 'FFFFFF':
                                has_white_text = True
                        except Exception:
                            pass
                        break
                break

            if is_teal and has_white_text:
                header_compound_pass += 1

        if header_compound_pass == num_cols:
            print(f"PASS: Component 2 — All {num_cols} header cells have teal bg + white text (0.15 pts)")
            total_score += 0.15
        elif header_compound_pass > 0:
            partial = 0.15 * (header_compound_pass / num_cols)
            print(f"PARTIAL: Component 2 — {header_compound_pass}/{num_cols} header cells have teal bg + white text ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No header cells have both teal bg and white text")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ===== Component 3: Even data rows (2,4,6) have light teal #E0F2F1 (0.30 pts) =====
    # Initial state has all rows white; only the teal stripe rows are task-introduced
    try:
        teal_rows_pass = 0
        teal_row_indices = [2, 4, 6]

        for r in teal_row_indices:
            row_ok = True
            for c in range(num_cols):
                fill = get_cell_fill_color(table.cell(r, c))
                if fill != 'E0F2F1':
                    row_ok = False
                    print(f"  DEBUG: Row {r}, Col {c} fill={fill}, expected E0F2F1")
                    break
            if row_ok:
                teal_rows_pass += 1

        if teal_rows_pass == len(teal_row_indices):
            print(f"PASS: Component 3 — All 3 even data rows have light teal (#E0F2F1) background (0.30 pts)")
            total_score += 0.30
        elif teal_rows_pass > 0:
            partial = 0.30 * (teal_rows_pass / len(teal_row_indices))
            print(f"PARTIAL: Component 3 — {teal_rows_pass}/{len(teal_row_indices)} even data rows have light teal ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No even data rows have light teal (#E0F2F1) background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ===== Component 4: Outer borders removed — noFill (0.25 pts) =====
    # Initial state has solid grey borders on all edges; removing outer borders is task-introduced
    try:
        outer_border_ok = 0
        outer_border_total = 0

        for r in range(num_rows):
            for c in range(num_cols):
                cell = table.cell(r, c)
                if r == 0:  # top edge
                    outer_border_total += 1
                    if get_border_info(cell, 'lnT') == 'noFill':
                        outer_border_ok += 1
                if r == num_rows - 1:  # bottom edge
                    outer_border_total += 1
                    if get_border_info(cell, 'lnB') == 'noFill':
                        outer_border_ok += 1
                if c == 0:  # left edge
                    outer_border_total += 1
                    if get_border_info(cell, 'lnL') == 'noFill':
                        outer_border_ok += 1
                if c == num_cols - 1:  # right edge
                    outer_border_total += 1
                    if get_border_info(cell, 'lnR') == 'noFill':
                        outer_border_ok += 1

        outer_ratio = outer_border_ok / outer_border_total if outer_border_total > 0 else 0

        if outer_ratio >= 0.9:
            print(f"PASS: Component 4 — Outer borders removed ({outer_border_ok}/{outer_border_total}) (0.25 pts)")
            total_score += 0.25
        elif outer_ratio > 0:
            partial = 0.25 * outer_ratio
            print(f"PARTIAL: Component 4 — Outer borders: {outer_border_ok}/{outer_border_total} removed ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No outer borders removed ({outer_border_ok}/{outer_border_total})")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
