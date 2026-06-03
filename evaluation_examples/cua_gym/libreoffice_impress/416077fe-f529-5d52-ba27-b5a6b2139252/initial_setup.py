"""
Initial Setup: Statistics presentation with plain table on slide 5
Task ID: impress_gf1_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_cell_border(cell, color_hex="808080", width_pt=0.5):
    """Set all 4 borders of a table cell to a thin solid line."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    width_emu = int(width_pt * 12700)  # 1pt = 12700 EMU

    for border_tag in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
        ln = tcPr.makeelement(qn(border_tag), {'w': str(width_emu), 'cap': 'flat', 'cmpd': 'sng'})
        solidFill = ln.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        prstDash = ln.makeelement(qn('a:prstDash'), {'val': 'solid'})
        ln.append(prstDash)
        # Remove existing border of same type
        existing = tcPr.find(qn(border_tag))
        if existing is not None:
            tcPr.remove(existing)
        tcPr.append(ln)


def set_cell_fill(cell, color_hex):
    """Set solid fill color for a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing solidFill
    for sf in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(sf)
    # Remove noFill if present
    for nf in tcPr.findall(qn('a:noFill')):
        tcPr.remove(nf)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Title Slide ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Business Statistics"
    slide1.placeholders[1].text = "Q1 2025 Performance Review\nPrepared by Analytics Division"

    # === Slide 2: Executive Summary ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Revenue grew 12.3% year-over-year, reaching $4.87M in Q1 2025."
    p2 = tf2.add_paragraph()
    p2.text = "Customer acquisition costs decreased by 8.1% compared to Q4 2024."
    p3 = tf2.add_paragraph()
    p3.text = "Employee satisfaction scores averaged 4.2/5.0 across all departments."
    p4 = tf2.add_paragraph()
    p4.text = "Market share in core segments increased to 23.7%, up from 21.4%."

    # === Slide 3: Revenue Breakdown ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Breakdown by Region"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "North America: $2.14M (43.9%)"
    for line in ["Europe: $1.32M (27.1%)", "Asia-Pacific: $0.98M (20.1%)", "Latin America: $0.43M (8.9%)"]:
        p = tf3.add_paragraph()
        p.text = line

    # === Slide 4: Key Metrics ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Performance Indicators"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Monthly Active Users: 1.24M (+15.6%)"
    for line in [
        "Average Order Value: $67.30 (+3.2%)",
        "Customer Retention Rate: 89.4% (+1.8pp)",
        "Net Promoter Score: 72 (+5 points)",
        "Support Ticket Resolution: 4.2 hours (-12%)",
    ]:
        p = tf4.add_paragraph()
        p.text = line

    # === Slide 5: Data Table (plain formatting) ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title text box
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.7))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Department Performance Summary"
    title_run = title_p.runs[0]
    title_run.font.size = Pt(24)
    title_run.font.bold = True

    # 5-column, 7-row table
    rows, cols = 7, 5
    table_shape = slide5.shapes.add_table(
        rows, cols,
        Inches(0.8), Inches(1.2),
        Inches(11.5), Inches(4.8)
    )
    table = table_shape.table

    # Headers
    headers = ["Department", "Revenue ($K)", "Headcount", "Growth (%)", "Satisfaction"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(14)

    # Data rows - realistic department data
    data = [
        ["Engineering", "1,245.60", "87", "14.2", "4.3"],
        ["Sales", "982.30", "64", "11.8", "3.9"],
        ["Marketing", "634.50", "42", "9.5", "4.1"],
        ["Operations", "512.80", "53", "7.3", "4.0"],
        ["Finance", "389.20", "31", "5.1", "4.4"],
        ["Human Resources", "278.40", "24", "3.8", "4.6"],
    ]

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                if c >= 1:  # right-align numeric columns
                    para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(12)

    # Apply plain uniform white fill and thin borders to ALL cells
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            set_cell_fill(cell, "FFFFFF")
            set_cell_border(cell, color_hex="808080", width_pt=0.5)

    # === Slide 6: Trends ===
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Emerging Trends & Outlook"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "AI-driven automation projects expected to save $320K annually by Q3."
    for line in [
        "Remote workforce productivity up 6.2% with hybrid model adoption.",
        "Sustainability initiatives on track to reduce carbon footprint by 15%.",
        "New product pipeline includes 3 launches scheduled for Q2 2025.",
    ]:
        p = tf6.add_paragraph()
        p.text = line

    # === Slide 7: Next Steps ===
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Next Steps & Action Items"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Finalize Q2 budget allocations by April 15, 2025."
    for line in [
        "Launch customer feedback survey across all regions.",
        "Complete hiring for 12 open engineering positions.",
        "Present board update with revised annual projections.",
        "Schedule department-level strategy sessions for May.",
    ]:
        p = tf7.add_paragraph()
        p.text = line

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
