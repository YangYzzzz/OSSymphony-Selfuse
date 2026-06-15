"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert the current date (fixed) under Heading 1 'Revision', right-aligned.
Generated: 2025-10-17 15:13:41
Status: success
Model: azure-o3
Total Steps: 3
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import os
import re
import datetime

FILE_PATH = "/home/user/insert_the_current_date_fixed_under_heading_1_revision_right_aligned.pptx"

def is_date_string(text: str) -> bool:
    """Return True iff *text* looks like a valid, human-readable date string."""
    text = text.strip()

    # Try to parse against several common date formats
    candidate_formats = [
        "%B %d, %Y", "%B %d %Y",          # October 14, 2025 | October 14 2025
        "%d %B %Y", "%d %B, %Y",          # 14 October 2025 | 14 October, 2025
        "%Y-%m-%d", "%m/%d/%Y",           # 2025-10-14 | 10/14/2025
    ]
    for fmt in candidate_formats:
        try:
            datetime.datetime.strptime(text, fmt)
            return True
        except ValueError:
            continue

    # Fallback regex: <MonthName> <day>[,] <year>
    month_names = (
        "January|February|March|April|May|June|July|August|"
        "September|October|November|December"
    )
    pattern = rf"^(?:{month_names})\s+\d{{1,2}},?\s+\d{{4}}$"
    return re.match(pattern, text) is not None


def verify_task(file_path: str) -> float:
    """Verify the task requirements and return a progressive score (0.0-1.0)."""
    print(f"Verifying presentation: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0

    # 1) Load the presentation -------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slide(s)")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    heading_shape = None  # The shape containing the heading text 'Revision'
    date_shape = None     # The shape containing the date string

    # 2) Locate heading "Revision" and date on the same slide ------------------
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            if (shape.text or "").strip().lower() == "revision":
                heading_shape = shape
                heading_slide = slide
                print(f"✓ Found heading 'Revision' on slide {slide_idx}")
                break
        if heading_shape:
            # Search for a date string in a different shape on the same slide
            for candidate in heading_slide.shapes:
                if candidate is heading_shape:
                    continue
                if hasattr(candidate, "text") and candidate.text and is_date_string(candidate.text):
                    date_shape = candidate
                    print(f"✓ Found date text '{candidate.text.strip()}' on the same slide")
                    break
            break  # Stop searching other slides once found

    # ----------------------- Progressive Scoring ------------------------------
    score = 0.0

    # Requirement A: Heading exists (0.4)
    if heading_shape:
        score += 0.4
    else:
        print("✗ Heading 'Revision' not found")
        print(f"REWARD: {score}")
        return score  # Cannot evaluate further without the heading

    # Requirement B: Date string exists (0.3)
    if date_shape:
        score += 0.3
    else:
        print("✗ Date text not found on the same slide as heading")
        print(f"REWARD: {score}")
        return score

    # Requirement C: Date is right-aligned (0.2)
    right_aligned = False
    if date_shape.has_text_frame:
        for paragraph in date_shape.text_frame.paragraphs:
            if paragraph.alignment in (PP_ALIGN.RIGHT, 3, "RIGHT"):
                right_aligned = True
                break
    if right_aligned:
        print("✓ Date text is right-aligned")
        score += 0.2
    else:
        print("✗ Date text is not right-aligned")

    # Requirement D: Date positioned below the heading (0.1)
    try:
        if date_shape.top > heading_shape.top:
            print("✓ Date text is positioned below the heading")
            score += 0.1
        else:
            print("✗ Date text is not positioned below the heading")
    except Exception as e:
        print(f"! Could not compare shape positions: {e}")

    # Correct any floating-point rounding artefacts
    if score > 0.999:
        score = 1.0

    print(f"Total score: {score}")
    return score

# ---------------------------- EXECUTION --------------------------------------
if __name__ == "__main__":
    reward = verify_task(FILE_PATH)
    print(f"REWARD: {reward}")
