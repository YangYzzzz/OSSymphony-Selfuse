"""
FINAL REWARD SCRIPT - SUCCESS
Task: Convert the semicolon-separated list in paragraph 4 to a table using ';' as delimiter.
Generated: 2025-10-17 09:05:57
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation

# -----------------------------------------------------------
# Reward Script for:
# "Convert the semicolon-separated list in paragraph 4 to a
#  table using ';' as delimiter."
# -----------------------------------------------------------
# This script verifies that the user has:
#   1. Inserted (at least) one table into the presentation
#   2. Populated that table with the expected items that were
#      originally separated by semicolons
#   3. Removed the original semicolon-separated paragraph so
#      no stray ';' characters remain in visible slide text
#
# Scoring (progressive – max 1.0):
#   • Table present ..................................... 0.4
#   • Each expected token found in table ................ 0.4
#       (0.4 ÷ #tokens each)
#   • No remaining semicolon list in slide text ......... 0.2
#
# A perfect score of 1.0 is awarded only when ALL checks pass.
# -----------------------------------------------------------

FILE_PATH = "/home/user/convert_the_semicolon_separated_list_in_paragraph_4_to_a_table_using_as_delimiter.pptx"

# Expected items that should now appear as individual cells
EXPECTED_TOKENS = [
    "Apples",
    "Oranges",
    "Bananas",
    "Grapes"
]


def verify_task(file_path: str) -> float:
    """Return a reward score between 0.0 and 1.0 based on task completion."""
    print(f"Verifying presentation for task: {file_path}\n")
    score = 0.0  # progressive score accumulator
    MAX_SCORE = 1.0

    # ------------------------------------------------------------------
    # 1) Load presentation (no points – prerequisite)
    # ------------------------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slide(s)")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0  # cannot continue without the file

    # ------------------------------------------------------------------
    # 2) Detect tables (0.4 points if at least one exists)
    # ------------------------------------------------------------------
    tables = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                tables.append((slide_idx, shape.table))
    if tables:
        score += 0.4
        print(f"✓ Found {len(tables)} table(s) in the presentation (+0.4)")
    else:
        print("✗ No tables found (0 points)")

    # ------------------------------------------------------------------
    # 3) Verify each expected token appears in some table cell
    #    Distribute 0.4 points equally across all tokens
    # ------------------------------------------------------------------
    token_points = 0.0
    if tables:
        # Gather ALL texts from ALL tables into a set for quick lookup
        table_texts = set()
        for _slide_idx, tbl in tables:
            for r in range(len(tbl.rows)):
                for c in range(len(tbl.columns)):
                    txt = tbl.cell(r, c).text_frame.text.strip()
                    if txt:
                        table_texts.add(txt)

        per_token_value = 0.4 / len(EXPECTED_TOKENS)
        for token in EXPECTED_TOKENS:
            if token in table_texts:
                token_points += per_token_value
                print(f"✓ Token '{token}' found in table (+{per_token_value:.2f})")
            else:
                print(f"✗ Token '{token}' NOT found in any table")
    else:
        print("Skipping token verification because no table was detected.")

    score += token_points

    # ------------------------------------------------------------------
    # 4) Ensure no remaining semicolon-separated text (0.2 points)
    # ------------------------------------------------------------------
    semicolon_text_present = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and ";" in shape.text:
                semicolon_text_present = True
                print(f"Found semicolon text still present: '{shape.text[:50]}...' (0 points)")
                break
        if semicolon_text_present:
            break

    if not semicolon_text_present:
        score += 0.2
        print("✓ No text with semicolons present (+0.2)")

    # ------------------------------------------------------------------
    # Final score
    # ------------------------------------------------------------------
    final_score = round(min(score, MAX_SCORE), 2)
    print(f"\nTotal score: {final_score}/{MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score


# -----------------------------------------------------------
# Execute verification when script is run directly
# -----------------------------------------------------------
if __name__ == "__main__":
    verify_task(FILE_PATH)

