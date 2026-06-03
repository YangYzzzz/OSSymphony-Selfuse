"""
Initial Setup: Reorder slides so slides 8-10 become slides 1-3
Task ID: impress_gf3_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === SLIDES 1-7: Main body content ===

    # Slide 1: Market Analysis Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide1.shapes.title.text = "Market Analysis Overview"
    add_textbox(slide1, 0.8, 1.8, 11.5, 1.0,
                "Q4 2025 saw a 12% increase in market share across all product lines.",
                font_size=20)
    add_textbox(slide1, 0.8, 3.0, 11.5, 1.0,
                "Key growth drivers: Enterprise SaaS adoption (+18%), SMB expansion (+9%), APAC penetration (+22%)",
                font_size=16)
    add_textbox(slide1, 0.8, 4.2, 11.5, 1.0,
                "Competitive landscape remains favorable with two major competitors exiting the mid-market segment.",
                font_size=16)

    # Slide 2: Revenue Breakdown
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Revenue Breakdown by Segment"
    add_textbox(slide2, 0.8, 1.8, 5.5, 3.5,
                "Enterprise: $45.2M (+15% YoY)\nMid-Market: $28.7M (+11% YoY)\nSMB: $12.3M (+8% YoY)\nChannel Partners: $8.9M (+21% YoY)",
                font_size=18)
    add_textbox(slide2, 6.8, 1.8, 5.5, 3.5,
                "Total Revenue: $95.1M\nGross Margin: 72.4%\nNet Revenue Retention: 118%\nNew ARR: $14.6M",
                font_size=18, bold=True)

    # Slide 3: Product Roadmap
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Product Roadmap - H1 2026"
    add_textbox(slide3, 0.8, 1.8, 11.5, 1.0,
                "Phase 1 (Jan-Feb): AI-powered analytics dashboard with predictive modeling",
                font_size=16)
    add_textbox(slide3, 0.8, 2.8, 11.5, 1.0,
                "Phase 2 (Mar-Apr): Multi-tenant architecture migration and performance optimization",
                font_size=16)
    add_textbox(slide3, 0.8, 3.8, 11.5, 1.0,
                "Phase 3 (May-Jun): Enterprise SSO integration and compliance certifications (SOC2, HIPAA)",
                font_size=16)
    add_textbox(slide3, 0.8, 4.8, 11.5, 1.0,
                "Budget allocation: $3.2M engineering, $1.1M infrastructure, $0.8M QA",
                font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # Slide 4: Customer Success Metrics
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "Customer Success Metrics"
    add_textbox(slide4, 0.8, 1.8, 11.5, 1.0,
                "NPS Score: 72 (up from 65 in Q3) | CSAT: 4.6/5.0 | Churn Rate: 2.1%",
                font_size=18, bold=True)
    add_textbox(slide4, 0.8, 3.0, 11.5, 1.0,
                "Top accounts renewed: Meridian Corp ($2.4M), Atlas Industries ($1.8M), Pinnacle Health ($1.5M)",
                font_size=16)
    add_textbox(slide4, 0.8, 4.2, 11.5, 1.0,
                "Support ticket resolution: Average 4.2 hours (target: <6 hours). 94% first-contact resolution.",
                font_size=16)

    # Slide 5: Team Growth & Hiring
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "Team Growth & Organizational Updates"
    add_textbox(slide5, 0.8, 1.8, 11.5, 1.0,
                "Current headcount: 342 (up from 298 in Q3 2025)",
                font_size=18)
    add_textbox(slide5, 0.8, 2.8, 5.5, 2.0,
                "Engineering: 156 (+18)\nSales: 67 (+8)\nCustomer Success: 45 (+6)\nMarketing: 34 (+5)",
                font_size=16)
    add_textbox(slide5, 6.8, 2.8, 5.5, 2.0,
                "Key hires:\n- VP Engineering: Priya Sharma\n- Director of Sales EMEA: Thomas Mueller\n- Head of Data Science: Kenji Tanaka",
                font_size=16)

    # Slide 6: Competitive Landscape
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    slide6.shapes.title.text = "Competitive Landscape Analysis"
    add_textbox(slide6, 0.8, 1.8, 11.5, 1.0,
                "Win rate vs. competitors: Nexus Pro (68%), CloudBase (72%), DataSync (81%)",
                font_size=18, bold=True)
    add_textbox(slide6, 0.8, 3.0, 11.5, 1.0,
                "Key differentiators: Real-time collaboration, AI-driven insights, seamless API integrations",
                font_size=16)
    add_textbox(slide6, 0.8, 4.2, 11.5, 1.0,
                "Market positioning: Leader in Forrester Wave Q4 2025, Strong Performer in Gartner MQ",
                font_size=16)

    # Slide 7: Financial Outlook
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    slide7.shapes.title.text = "Financial Outlook - FY2026"
    add_textbox(slide7, 0.8, 1.8, 11.5, 1.0,
                "Revenue target: $420M (35% YoY growth) | EBITDA target: $84M (20% margin)",
                font_size=18, bold=True)
    add_textbox(slide7, 0.8, 3.0, 11.5, 1.0,
                "Key investments: International expansion ($15M), R&D acceleration ($22M), Sales capacity ($12M)",
                font_size=16)
    add_textbox(slide7, 0.8, 4.2, 11.5, 1.0,
                "Path to profitability: Targeting cash-flow positive by Q3 2026 with current growth trajectory.",
                font_size=16)

    # === SLIDES 8-10: Introduction/Context slides (should appear first after reorder) ===

    # Slide 8: Company Introduction
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    slide8.shapes.title.text = "Stratosphere Analytics - Company Overview"
    add_textbox(slide8, 0.8, 1.8, 11.5, 1.0,
                "Founded 2019 | Headquarters: San Francisco, CA | Series D ($180M raised)",
                font_size=18)
    add_textbox(slide8, 0.8, 3.0, 11.5, 1.0,
                "Mission: Empowering data-driven decisions through intelligent analytics.",
                font_size=20, bold=True, color=RGBColor(0x00, 0x70, 0xC0))
    add_textbox(slide8, 0.8, 4.5, 11.5, 1.0,
                "Serving 2,400+ customers across 38 countries in technology, healthcare, and financial services.",
                font_size=16)

    # Slide 9: Meeting Agenda
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    slide9.shapes.title.text = "Board Meeting Agenda - January 2026"
    add_textbox(slide9, 0.8, 1.8, 11.5, 4.0,
                "1. Company Overview & Mission Update\n"
                "2. Market Analysis & Competitive Position\n"
                "3. Revenue Performance & Financial Metrics\n"
                "4. Product Roadmap & Engineering Priorities\n"
                "5. Customer Success & Retention\n"
                "6. Team Growth & Key Hires\n"
                "7. Competitive Landscape\n"
                "8. Financial Outlook & FY2026 Targets",
                font_size=18)

    # Slide 10: Executive Summary
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    slide10.shapes.title.text = "Executive Summary"
    add_textbox(slide10, 0.8, 1.8, 11.5, 1.0,
                "Q4 2025 delivered record results: $95.1M revenue, 72 NPS, 118% net retention.",
                font_size=18, bold=True)
    add_textbox(slide10, 0.8, 3.0, 11.5, 1.0,
                "Strategic priorities for 2026: AI product innovation, international expansion, path to profitability.",
                font_size=18)
    add_textbox(slide10, 0.8, 4.2, 11.5, 1.0,
                "This presentation covers our market position, financial performance, and roadmap for the coming year.",
                font_size=16, color=RGBColor(0x66, 0x66, 0x66))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
