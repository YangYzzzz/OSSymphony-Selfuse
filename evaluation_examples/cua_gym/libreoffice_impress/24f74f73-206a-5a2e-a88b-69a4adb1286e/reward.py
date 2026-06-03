"""
Reward Script: Reorder slides so slides 8-10 become slides 1-3, original 1-7 shift to 4-10.
Task ID: impress_gf3_009
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Positions 1-3 contain former slides 8,9,10 (by title text)
  Component 2 (0.4): Positions 4-10 contain former slides 1-7 (by title text)
  Component 3 (0.2): Spot-check content integrity on reordered slides
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_009'

# Expected slide titles in the GOLDEN (reordered) presentation
# Former slide 8 -> new pos 1, former slide 9 -> new pos 2, former slide 10 -> new pos 3
# Former slides 1-7 -> new pos 4-10
EXPECTED_TITLES = [
    "Stratosphere Analytics - Company Overview",   # former slide 8 -> pos 1
    "Board Meeting Agenda - January 2026",          # former slide 9 -> pos 2
    "Executive Summary",                            # former slide 10 -> pos 3
    "Market Analysis Overview",                     # former slide 1 -> pos 4
    "Revenue Breakdown by Segment",                 # former slide 2 -> pos 5
    "Product Roadmap - H1 2026",                    # former slide 3 -> pos 6
    "Customer Success Metrics",                     # former slide 4 -> pos 7
    "Team Growth & Organizational Updates",         # former slide 5 -> pos 8
    "Competitive Landscape Analysis",               # former slide 6 -> pos 9
    "Financial Outlook - FY2026",                   # former slide 7 -> pos 10
]

# Content spot-checks: (slide_position_0indexed, substring_to_find)
# These verify content survived the reorder intact
CONTENT_CHECKS = [
    (0, "Founded 2019"),               # former slide 8 content at new pos 1
    (2, "$95.1M revenue"),             # former slide 10 content at new pos 3
    (3, "12% increase in market share"), # former slide 1 content at new pos 4
    (9, "$420M"),                       # former slide 7 content at new pos 10
]


def get_slide_title(slide):
    """Extract the title text from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return ""


def get_all_text(slide):
    """Get concatenated text from all shapes on a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return " ".join(texts)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: must have exactly 10 slides
    if len(prs.slides) != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Positions 1-3 have correct titles (former slides 8,9,10) — 0.4 points
    try:
        match_count = 0
        for i in range(3):
            actual_title = get_slide_title(slides[i])
            expected = EXPECTED_TITLES[i]
            if actual_title == expected:
                match_count += 1
                print(f"PASS: Slide {i+1} title matches: '{expected}'")
            else:
                print(f"FAIL: Slide {i+1} title mismatch: expected '{expected}', got '{actual_title}'")

        if match_count == 3:
            print(f"PASS: Component 1 — All 3 opening slides correct (0.4 pts)")
            total_score += 0.4
        elif match_count > 0:
            partial = round(0.4 * match_count / 3, 2)
            print(f"PARTIAL: Component 1 — {match_count}/3 opening slides correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No opening slides match expected order")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Positions 4-10 have correct titles (former slides 1-7) — 0.4 points
    try:
        match_count = 0
        for i in range(3, 10):
            actual_title = get_slide_title(slides[i])
            expected = EXPECTED_TITLES[i]
            if actual_title == expected:
                match_count += 1
                print(f"PASS: Slide {i+1} title matches: '{expected}'")
            else:
                print(f"FAIL: Slide {i+1} title mismatch: expected '{expected}', got '{actual_title}'")

        if match_count == 7:
            print(f"PASS: Component 2 — All 7 body slides in correct positions (0.4 pts)")
            total_score += 0.4
        elif match_count > 0:
            partial = round(0.4 * match_count / 7, 2)
            print(f"PARTIAL: Component 2 — {match_count}/7 body slides correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — No body slides match expected order")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Content integrity spot-checks (0.2 points)
    try:
        content_pass = 0
        for (idx, substr) in CONTENT_CHECKS:
            all_text = get_all_text(slides[idx])
            if substr in all_text:
                content_pass += 1
                print(f"PASS: Slide {idx+1} contains expected content '{substr}'")
            else:
                print(f"FAIL: Slide {idx+1} missing expected content '{substr}'")

        if content_pass == len(CONTENT_CHECKS):
            print(f"PASS: Component 3 — All content integrity checks passed (0.2 pts)")
            total_score += 0.2
        elif content_pass > 0:
            partial = round(0.2 * content_pass / len(CONTENT_CHECKS), 2)
            print(f"PARTIAL: Component 3 — {content_pass}/{len(CONTENT_CHECKS)} content checks passed ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No content integrity checks passed")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
