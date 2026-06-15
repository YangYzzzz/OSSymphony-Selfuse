"""
FINAL REWARD SCRIPT - SUCCESS
Task: Sort Table 1 by column 2 in ascending order as text.
Generated: 2025-10-17 06:20:05
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation


def verify_sort_table_task(file_path: str) -> float:
    """Verify that Table 1 in the given PPTX is sorted by column 2 in ascending
    order (as text).  Returns a progressive score between 0.0 and 1.0.
    """
    print(f"Verifying task on file: {file_path}")
    total_score = 0.0        # progressive score accumulator
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1. Load the presentation ------------------------------------------------
    # ------------------------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0  # cannot continue without loading

    # ------------------------------------------------------------------
    # 2. Locate the FIRST table ("Table 1") -----------------------------------
    # ------------------------------------------------------------------
    first_table = None
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                first_table = shape.table
                print(f"✓ Found a table on slide {slide_idx + 1}")
                break
        if first_table is not None:
            break

    if first_table is None:
        print("✗ No table found in presentation – task not completed")
        return 0.0

    # ------------------------------------------------------------------
    # 3. Basic table sanity checks --------------------------------------------
    # ------------------------------------------------------------------
    row_count = len(first_table.rows)
    col_count = len(first_table.columns)
    print(f"Table size: {row_count} rows × {col_count} columns")

    # Require at least a header + 2 data rows for meaningful sort check
    if row_count >= 3:
        total_score += 0.3  # 0.3 pts for having a data-sized table
        print("✓ Table has sufficient data rows (0.3 points)")
    else:
        print("✗ Table does not have enough data rows (need ≥ 3 rows)")
        # Continue – still possible to earn some points if sort is right

    if col_count < 2:
        print("✗ Table does not have at least 2 columns to sort by – verification stops")
        return total_score  # cannot assess sort without column 2

    # ------------------------------------------------------------------
    # 4. Extract data rows & verify sorting -----------------------------------
    # ------------------------------------------------------------------
    sort_col_idx = 1  # column 2 (0-based index)

    # Collect text from each data row (skip header row at index 0)
    data_rows = [
        [first_table.cell(r, c).text_frame.text.strip() for c in range(col_count)]
        for r in range(1, row_count)
    ]

    # Extract the key column values in current order
    key_values = [row[sort_col_idx] for row in data_rows]
    print("Column 2 values in data rows:", key_values)

    # Expected ascending order (case-insensitive textual sort)
    expected_order = sorted(key_values, key=lambda x: x.lower())

    if key_values == expected_order:
        total_score += 0.7  # 0.7 pts for correct sort
        print("✓ Data rows are sorted in ascending order as text by column 2 (0.7 points)")
    else:
        print("✗ Data rows are NOT sorted correctly by column 2")
        print("  Expected order:", expected_order)

    # ------------------------------------------------------------------
    # 5. Final score -----------------------------------------------------------
    # ------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"Final score: {final_score}")
    return final_score


if __name__ == "__main__":
    # Path provided by the task description
    FILE_PATH = "/home/user/sort_table_1_by_column_2_in_ascending_order_as_text.pptx"

    reward = verify_sort_table_task(FILE_PATH)
    print(f"REWARD: {reward}")
