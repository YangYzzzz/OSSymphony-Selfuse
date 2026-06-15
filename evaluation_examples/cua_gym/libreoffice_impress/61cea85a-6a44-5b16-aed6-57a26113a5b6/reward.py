"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert a non-breaking space after 'Mr.', 'Ms.' and 'Dr.' in paragraph 3.
Generated: 2025-10-17 11:39:52
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
import re
from pptx import Presentation

def verify_non_breaking_space_after_abbr(file_path: str) -> float:
    """Reward script for the task:
    "Insert a non-breaking space after 'Mr.', 'Ms.' and 'Dr.' in paragraph 3."

    Scoring (progressive):
        • 1/3 point for each abbreviation that is ONLY followed by a non-breaking
          space (\u00A0) somewhere in the presentation and never by a normal
          breaking space.  
        • Half-credit (1/6) if both NBSP and regular space are found after the
          same abbreviation (mixed formatting).
        • 0 if no NBSP detected.

    Returns a float between 0.0 and 1.0 and prints debug information plus the
    final score as "REWARD: X.X".
    """

    abbreviations = ["Mr.", "Ms.", "Dr."]

    # ------------------------------------------------------------------
    # 0. Preliminary checks
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slide(s)")
    except Exception as exc:
        print(f"✗ Failed to open presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 1. Scan text and count occurrences of NBSP vs normal spaces
    # ------------------------------------------------------------------
    nb_counts = {abbr: 0 for abbr in abbreviations}  # non-breaking space counts
    sp_counts = {abbr: 0 for abbr in abbreviations}  # normal space counts

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs, start=1):
                text = para.text or ""
                if not text:
                    continue
                for abbr in abbreviations:
                    nb_counts[abbr] += text.count(abbr + "\u00A0")
                    sp_counts[abbr] += len(re.findall(re.escape(abbr) + r" ", text))

    # ------------------------------------------------------------------
    # 2. Progressive scoring
    # ------------------------------------------------------------------
    total_score = 0.0
    max_per_abbr = 1.0 / len(abbreviations)

    print("\nVerification results:")
    for abbr in abbreviations:
        nb = nb_counts[abbr]
        sp = sp_counts[abbr]
        print(f"  {abbr}: NBSP={nb}, regular space={sp}")

        if nb > 0 and sp == 0:
            print(f"    ✓ Correct: only non-breaking spaces found after '{abbr}'")
            total_score += max_per_abbr
        elif nb > 0 and sp > 0:
            print(f"    ⚠ Mixed formatting after '{abbr}' (partial credit)")
            total_score += max_per_abbr * 0.5
        else:
            print(f"    ✗ No non-breaking space found after '{abbr}'")

    total_score = round(min(total_score, 1.0), 4)
    print(f"\nREWARD: {total_score}")
    return total_score

# ----------------------------------------------------------------------
# Execute verification when run directly
# ----------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_a_non_breaking_space_after_mr_ms_and_dr_in_paragraph_3.pptx"
    verify_non_breaking_space_after_abbr(FILE_PATH)
