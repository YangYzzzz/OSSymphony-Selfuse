"""
Reward Script: Create a data-driven Sales Report presentation from Sales_Data.xlsx
Task ID: impress_wf_013
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count == 7 (0.10)
  Component 2: Slide 1 title is 'Sales Report Q3' (0.10)
  Component 3: Slide 2 has summary numbers (revenue, units, avg order) (0.15)
  Component 4: Slide 3 has a bar/column chart (0.15)
  Component 5: Slide 4 has a pie chart (0.15)
  Component 6: Slide 5 has a line chart (0.15)
  Component 7: Slide 6 has a table with 6 rows (header + 5 products) (0.10)
  Component 8: Slide 7 has 4 bullet/takeaway points (0.10)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_013'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slide count == 7 (0.10 points)
    try:
        num_slides = len(slides)
        if num_slides == 7:
            print(f"PASS: Component 1 -- Slide count is 7 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 -- Expected 7 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If fewer than 7 slides, some checks below will fail gracefully
    # but we still attempt them in case the ordering is different

    # Component 2: Slide 1 title is 'Sales Report Q3' (0.10 points)
    try:
        if len(slides) >= 1:
            slide1 = slides[0]
            title_text = ""
            for shape in slide1.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        # Take the first non-empty text shape as potential title
                        if not title_text:
                            title_text = txt
                # Also check placeholder title
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    if shape.placeholder_format.idx == 0:  # title placeholder
                        title_text = shape.text_frame.text.strip()
                        break

            if "sales report q3" in title_text.lower():
                print(f"PASS: Component 2 -- Slide 1 title is '{title_text}' (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Expected title containing 'Sales Report Q3', found '{title_text}'")
        else:
            print(f"FAIL: Component 2 -- No slides found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 2 has summary numbers (0.15 points)
    # Check for presence of revenue, units, and avg order value styled numbers
    try:
        if len(slides) >= 2:
            slide2 = slides[1]
            all_text = ""
            for shape in slide2.shapes:
                if shape.has_text_frame:
                    all_text += " " + shape.text_frame.text.strip()

            all_text_lower = all_text.lower()
            has_revenue = "$" in all_text and any(
                kw in all_text_lower for kw in ["revenue", "total revenue"]
            )
            has_units = any(
                kw in all_text_lower for kw in ["unit", "units", "total units"]
            )
            has_avg = any(
                kw in all_text_lower for kw in ["avg", "average", "order value"]
            )

            checks_passed = sum([has_revenue, has_units, has_avg])
            if checks_passed == 3:
                print(f"PASS: Component 3 -- Slide 2 has revenue, units, and avg order value (0.15 pts)")
                total_score += 0.15
            elif checks_passed >= 2:
                print(f"PARTIAL: Component 3 -- Slide 2 has {checks_passed}/3 summary items (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 -- Slide 2 missing summary numbers (revenue={has_revenue}, units={has_units}, avg={has_avg})")
        else:
            print(f"FAIL: Component 3 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide 3 has a bar/column chart (0.15 points)
    try:
        if len(slides) >= 3:
            slide3 = slides[2]
            found_chart = False
            chart_type_val = None
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    found_chart = True
                    chart_type_val = shape.chart.chart_type
                    break

            if found_chart:
                # BAR types: BAR_CLUSTERED(57), BAR_STACKED(58), BAR_STACKED_100(59)
                # COLUMN types: COLUMN_CLUSTERED(51), COLUMN_STACKED(52), COLUMN_STACKED_100(53)
                # Accept any bar or column chart
                ct = int(chart_type_val)
                if ct in (51, 52, 53, 54, 57, 58, 59, 60):
                    print(f"PASS: Component 4 -- Slide 3 has bar/column chart (type={ct}) (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"PARTIAL: Component 4 -- Slide 3 has a chart but type={ct}, expected bar/column. (0.07 pts)")
                    total_score += 0.07
            else:
                print(f"FAIL: Component 4 -- Slide 3 has no chart")
        else:
            print(f"FAIL: Component 4 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Slide 4 has a pie chart (0.15 points)
    try:
        if len(slides) >= 4:
            slide4 = slides[3]
            found_chart = False
            chart_type_val = None
            for shape in slide4.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    found_chart = True
                    chart_type_val = shape.chart.chart_type
                    break

            if found_chart:
                # PIE types: PIE(5), PIE_EXPLODED(69), PIE_OF_PIE(68)
                ct = int(chart_type_val)
                if ct in (5, 68, 69, 70):
                    print(f"PASS: Component 5 -- Slide 4 has pie chart (type={ct}) (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"PARTIAL: Component 5 -- Slide 4 has a chart but type={ct}, expected pie. (0.07 pts)")
                    total_score += 0.07
            else:
                print(f"FAIL: Component 5 -- Slide 4 has no chart")
        else:
            print(f"FAIL: Component 5 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    # Component 6: Slide 5 has a line chart (0.15 points)
    try:
        if len(slides) >= 5:
            slide5 = slides[4]
            found_chart = False
            chart_type_val = None
            for shape in slide5.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    found_chart = True
                    chart_type_val = shape.chart.chart_type
                    break

            if found_chart:
                # LINE types: LINE(4), LINE_MARKERS(65), LINE_STACKED(63), etc.
                ct = int(chart_type_val)
                if ct in (4, 63, 64, 65, 66, 67):
                    print(f"PASS: Component 6 -- Slide 5 has line chart (type={ct}) (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"PARTIAL: Component 6 -- Slide 5 has a chart but type={ct}, expected line. (0.07 pts)")
                    total_score += 0.07
            else:
                print(f"FAIL: Component 6 -- Slide 5 has no chart")
        else:
            print(f"FAIL: Component 6 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 -- {e}")

    # Component 7: Slide 6 has a table with header + 5 data rows (0.10 points)
    try:
        if len(slides) >= 6:
            slide6 = slides[5]
            found_table = False
            table_rows = 0
            table_cols = 0
            for shape in slide6.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    found_table = True
                    table_rows = len(shape.table.rows)
                    table_cols = len(shape.table.columns)
                    break

            if found_table and table_rows >= 6:
                print(f"PASS: Component 7 -- Slide 6 has table with {table_rows} rows, {table_cols} cols (0.10 pts)")
                total_score += 0.10
            elif found_table and table_rows >= 2:
                print(f"PARTIAL: Component 7 -- Slide 6 has table but only {table_rows} rows (expected >=6) (0.05 pts)")
                total_score += 0.05
            elif found_table:
                print(f"FAIL: Component 7 -- Slide 6 has table but only {table_rows} rows")
            else:
                print(f"FAIL: Component 7 -- Slide 6 has no table")
        else:
            print(f"FAIL: Component 7 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 -- {e}")

    # Component 8: Slide 7 has 4 bullet/takeaway points (0.10 points)
    try:
        if len(slides) >= 7:
            slide7 = slides[6]
            # Collect all non-empty paragraphs from non-title text shapes
            bullet_texts = []
            for shape in slide7.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text_frame.text.strip().lower()
                    # Skip title-like shapes (contain "takeaway" or "key")
                    is_title_shape = False
                    paras = shape.text_frame.paragraphs
                    if len(paras) <= 1:
                        # Single paragraph shape -- could be a title
                        if any(kw in shape_text for kw in ["takeaway", "key takeaway", "key findings"]):
                            is_title_shape = True
                    if not is_title_shape:
                        for para in paras:
                            txt = para.text.strip()
                            if txt and "takeaway" not in txt.lower()[:15]:
                                bullet_texts.append(txt)

            num_bullets = len(bullet_texts)
            if num_bullets >= 4:
                print(f"PASS: Component 8 -- Slide 7 has {num_bullets} bullet points (0.10 pts)")
                total_score += 0.10
            elif num_bullets >= 2:
                print(f"PARTIAL: Component 8 -- Slide 7 has {num_bullets} bullet points, expected 4 (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 -- Slide 7 has {num_bullets} bullet points, expected 4")
        else:
            print(f"FAIL: Component 8 -- Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Desktop/Sales_Report.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
