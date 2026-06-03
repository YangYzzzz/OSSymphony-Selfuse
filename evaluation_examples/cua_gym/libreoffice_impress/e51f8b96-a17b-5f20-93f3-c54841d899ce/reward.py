"""
FINAL REWARD SCRIPT - SUCCESS
Task: Set the line spacing of paragraphs 2–3 to Exactly 18 pt.
Generated: 2025-10-17 11:00:15
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def _is_18pt(spacing_val, tolerance_pt: float = 1.0):
    """Return True if `spacing_val` (EMU units) is ~18 pt within tolerance."""
    if spacing_val is None:
        return False
    # 1 point (pt) = 12700 English Metric Units (EMU) inside PPTX XML
    target = 18 * 12700
    tolerance = tolerance_pt * 12700
    return abs(spacing_val - target) <= tolerance


def verify_line_spacing_18pt(file_path: str) -> float:
    """Verify that paragraphs 2–3 have Exactly 18 pt line spacing.

    Returns a progressive score between 0.0 and 1.0.
    Scoring rules:
        • 0.7 points if BOTH paragraph 2 and paragraph 3 in a text frame are
          set to ~18 pt (Exactly).
        • Additional 0.3 points if all other paragraphs in that text frame do
          NOT have 18 pt spacing (i.e., only paragraphs 2–3 were modified).
    """
    print(f"Verifying presentation: {file_path}")

    # PRECONDITION: File must exist and load successfully (no points awarded)
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Error opening PPTX: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # Flags for scoring
    para2_correct = False
    para3_correct = False
    other_paras_ok = True  # assume true until proven otherwise
    found_frame_with_target = False

    # Iterate over all slides & shapes to locate a text frame that contains
    # at least 3 paragraphs where paragraphs 2–3 meet the requirement.
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not getattr(shape, "has_text_frame", False):
                continue

            paragraphs = shape.text_frame.paragraphs
            if len(paragraphs) < 3:
                continue  # need at least 3 paragraphs to evaluate

            # Evaluate paragraph 2 (index 1) & paragraph 3 (index 2)
            para2_correct = _is_18pt(paragraphs[1].line_spacing)
            para3_correct = _is_18pt(paragraphs[2].line_spacing)

            # Early skip if either paragraph fails – we'll keep searching
            if not (para2_correct and para3_correct):
                continue

            # Check that OTHER paragraphs do NOT have 18 pt spacing.
            other_paras_ok = True
            for idx, para in enumerate(paragraphs):
                if idx in (1, 2):
                    continue  # Skip paragraphs 2 & 3 (already verified)
                if _is_18pt(para.line_spacing):
                    other_paras_ok = False
                    break

            found_frame_with_target = True
            break  # Stop searching shapes
        if found_frame_with_target:
            break  # Stop searching slides

    # Scoring according to verification results
    score = 0.0
    if para2_correct and para3_correct:
        print("✓ Paragraphs 2 and 3 have exactly 18 pt line spacing")
        score += 0.7
        if other_paras_ok:
            print("✓ Other paragraphs NOT set to 18 pt – correct scope of change")
            score += 0.3
        else:
            print("✗ Some other paragraphs also have 18 pt spacing (should remain unchanged)")
    elif para2_correct or para3_correct:
        # Partial success: only one paragraph correctly formatted
        correct_para_num = 2 if para2_correct else 3
        print(f"✓ Paragraph {correct_para_num} correctly set to 18 pt, other incorrect")
        score += 0.3
    else:
        print("✗ Paragraphs 2 and 3 do NOT have the required 18 pt spacing")

    # Cap score at 1.0 and print final result
    score = min(score, 1.0)
    print(f"REWARD: {score}")
    return score


# ---------------------------------------------------------------------------
# When executed as a stand-alone script, run the verification immediately.
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/set_the_line_spacing_of_paragraphs_23_to_exactly_18_pt.pptx"
    verify_line_spacing_18pt(FILE_PATH)
