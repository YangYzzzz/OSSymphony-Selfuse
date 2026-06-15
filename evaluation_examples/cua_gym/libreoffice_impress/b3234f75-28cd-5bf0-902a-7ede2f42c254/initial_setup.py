"""
Initial Setup: Create Tech_Architecture.pptx with 8 slides, slide 4 has title 'Request Flow' but empty content.
Task ID: impress_ps_030
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_030'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Set the title and populate body placeholder with bullet lines."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        para = tf.paragraphs[i]
        para.level = 0
        for run in para.runs:
            run.font.size = Pt(18)


def add_title_only(slide, title_text):
    """Set only the title on a slide."""
    slide.shapes.title.text = title_text
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.bold = True


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Layout references
    layout_title = prs.slide_layouts[0]       # Title Slide
    layout_content = prs.slide_layouts[1]     # Title and Content
    layout_blank = prs.slide_layouts[5]       # Blank
    layout_title_only = prs.slide_layouts[5]  # We'll use blank for slide 4

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(layout_title)
    slide1.shapes.title.text = "Tech Architecture Overview"
    slide1.placeholders[1].text = "Cloud-Native Microservices Platform\nQ2 2025 Architecture Review"

    # --- Slide 2: System Components ---
    slide2 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide2, "System Components", [
        "Frontend: React SPA with SSR via Next.js",
        "API Layer: Express.js gateway with rate limiting",
        "Auth: OAuth2 + JWT token management",
        "Cache: Redis cluster with 99.9% hit rate",
        "Message Queue: RabbitMQ for async processing",
        "Storage: PostgreSQL primary + read replicas",
    ])

    # --- Slide 3: Infrastructure Overview ---
    slide3 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide3, "Infrastructure Overview", [
        "AWS EKS clusters across 3 availability zones",
        "Terraform-managed infrastructure as code",
        "CloudFront CDN for static asset delivery",
        "Auto-scaling groups: min 3, max 12 instances",
        "Datadog monitoring with PagerDuty alerts",
        "Weekly disaster recovery drills",
    ])

    # --- Slide 4: Request Flow (EMPTY content area - just title) ---
    slide4 = prs.slides.add_slide(layout_blank)
    # Add only a title text box at the top
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Request Flow"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # NO other shapes or connectors on this slide

    # --- Slide 5: Performance Metrics ---
    slide5 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide5, "Performance Metrics", [
        "Average API response time: 45ms (p50), 120ms (p99)",
        "Database query latency: 8ms average",
        "Cache hit ratio: 94.7% across all endpoints",
        "Uptime SLA: 99.95% achieved last quarter",
        "Peak throughput: 12,400 requests/second",
        "Error rate: 0.02% (4xx + 5xx combined)",
    ])

    # --- Slide 6: Security Architecture ---
    slide6 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide6, "Security Architecture", [
        "TLS 1.3 enforced for all external traffic",
        "Mutual TLS between internal microservices",
        "AWS WAF rules with custom threat signatures",
        "Secrets managed via HashiCorp Vault",
        "SOC 2 Type II certification in progress",
        "Quarterly penetration testing schedule",
    ])

    # --- Slide 7: Deployment Pipeline ---
    slide7 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide7, "Deployment Pipeline", [
        "GitHub Actions CI with parallel test suites",
        "Docker multi-stage builds for minimal images",
        "Canary deployments: 5% → 25% → 100% rollout",
        "Automated rollback on error rate spike > 1%",
        "Blue-green deployments for database migrations",
        "Average deploy cycle: commit to production in 18 minutes",
    ])

    # --- Slide 8: Summary & Next Steps ---
    slide8 = prs.slides.add_slide(layout_content)
    add_title_and_body(slide8, "Summary & Next Steps", [
        "Architecture handles 10x current load headroom",
        "Migration to gRPC for inter-service communication (Q3)",
        "Evaluate GraphQL federation for API gateway (Q3)",
        "Implement distributed tracing with OpenTelemetry (Q4)",
        "Cost optimization review: target 15% infra savings",
        "Hire 2 additional SRE engineers by end of year",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
