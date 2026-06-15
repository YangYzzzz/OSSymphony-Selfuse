"""
Reward Script: Data flow diagram on slide 4 of Tech_Architecture.pptx
Task ID: impress_ps_030
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Six process rectangle shapes with correct labels
  Component 2 (0.20): Diamond decision shape labeled 'Authenticated?'
  Component 3 (0.20): Arrow connectors linking shapes (>= 5 connectors)
  Component 4 (0.15): 'Yes' and 'No' decision path labels
  Component 5 (0.10): Error Response shape for the 'No' path
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_030'

# Expected process labels (rectangles) that must be on slide 4
EXPECTED_PROCESS_LABELS = [
    'User Request',
    'Load Balancer',
    'API Gateway',
    'Auth Service',
    'Business Service',
    'Database',
]


def normalize_text(text):
    """Normalize text for comparison: strip, collapse whitespace, lowercase."""
    if not text:
        return ''
    return ' '.join(text.strip().lower().split())


def get_slide4_shapes(prs):
    """Get all shapes on slide 4 (index 3), excluding placeholders and the title text box."""
    if len(prs.slides) < 4:
        return []
    slide = prs.slides[3]
    shapes = []
    for shape in slide.shapes:
        # Skip placeholders (title, content placeholders)
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        shapes.append(shape)
    return shapes


def verify_task(file_path):
    """
    Verify data flow diagram creation on slide 4.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]
    shapes = get_slide4_shapes(prs)

    # Categorize shapes
    auto_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    line_shapes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    text_boxes = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]

    # Separate auto_shapes into rectangles and diamonds
    rectangles = []
    diamonds = []
    for s in auto_shapes:
        try:
            if s.auto_shape_type is not None and s.auto_shape_type == 4:  # DIAMOND
                diamonds.append(s)
            else:
                rectangles.append(s)
        except Exception:
            rectangles.append(s)

    print(f"INFO: Slide 4 has {len(auto_shapes)} auto_shapes, {len(line_shapes)} lines, {len(text_boxes)} text_boxes")
    print(f"INFO: {len(rectangles)} rectangles, {len(diamonds)} diamonds")

    # Get all text from text_boxes (for Yes/No label check)
    text_box_texts = [normalize_text(tb.text) for tb in text_boxes if hasattr(tb, 'text')]

    # Get all text from rectangles
    rect_texts = [normalize_text(r.text) for r in rectangles if hasattr(r, 'text')]

    # Component 1: Six process rectangle shapes with correct labels (0.35 points)
    # Each of the 6 expected labels found earns proportional credit
    try:
        found_labels = []
        for expected in EXPECTED_PROCESS_LABELS:
            expected_norm = normalize_text(expected)
            # Check if any rectangle contains this label
            matched = any(expected_norm in rt or rt in expected_norm for rt in rect_texts)
            if matched:
                found_labels.append(expected)

        label_ratio = len(found_labels) / len(EXPECTED_PROCESS_LABELS)
        comp1_score = round(0.35 * label_ratio, 4)

        if label_ratio == 1.0:
            print(f"PASS: Component 1 -- All 6 process rectangles found ({comp1_score} pts)")
        else:
            missing = [l for l in EXPECTED_PROCESS_LABELS if l not in found_labels]
            print(f"PARTIAL: Component 1 -- {len(found_labels)}/6 process rectangles found, missing: {missing} ({comp1_score} pts)")
        total_score += comp1_score
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Diamond decision shape labeled 'Authenticated?' (0.20 points)
    try:
        diamond_with_auth = any(
            'authenticated' in (normalize_text(d.text) if hasattr(d, 'text') else '')
            for d in diamonds
        )

        if diamond_with_auth:
            print(f"PASS: Component 2 -- Diamond shape with 'Authenticated?' label found (0.20 pts)")
            total_score += 0.20
        elif len(diamonds) > 0:
            # Diamond exists but wrong label
            d_texts = [d.text for d in diamonds if hasattr(d, 'text')]
            print(f"PARTIAL: Component 2 -- Diamond found but label mismatch: {d_texts} (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 -- No diamond shape found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Arrow connectors (>= 5 line connectors) (0.20 points)
    try:
        connector_count = len(line_shapes)
        if connector_count >= 5:
            print(f"PASS: Component 3 -- {connector_count} connectors found (>= 5 required) (0.20 pts)")
            total_score += 0.20
        elif connector_count >= 3:
            partial = round(0.20 * (connector_count / 5), 4)
            print(f"PARTIAL: Component 3 -- {connector_count} connectors found (>= 5 required) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Only {connector_count} connectors found (>= 5 required)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: 'Yes' and 'No' decision path labels (0.15 points)
    try:
        # Check text boxes and also auto_shapes for Yes/No labels
        all_texts = text_box_texts + rect_texts
        # Also check all shape texts on the slide (some labels may be in other shape types)
        for s in shapes:
            if hasattr(s, 'text') and s.text:
                all_texts.append(normalize_text(s.text))

        has_yes = any('yes' == t for t in all_texts)
        has_no = any('no' == t for t in all_texts)

        if has_yes and has_no:
            print(f"PASS: Component 4 -- Both 'Yes' and 'No' labels found (0.15 pts)")
            total_score += 0.15
        elif has_yes or has_no:
            found = 'Yes' if has_yes else 'No'
            print(f"PARTIAL: Component 4 -- Only '{found}' label found (0.075 pts)")
            total_score += 0.075
        else:
            print(f"FAIL: Component 4 -- Neither 'Yes' nor 'No' labels found")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Error Response shape for 'No' path (0.10 points)
    try:
        has_error_shape = any(
            'error' in (normalize_text(s.text) if hasattr(s, 'text') else '')
            for s in auto_shapes
        )

        if has_error_shape:
            print(f"PASS: Component 5 -- Error Response shape found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 -- No Error Response shape found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
