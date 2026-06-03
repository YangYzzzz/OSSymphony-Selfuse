"""
Reward Script: Build organizational chart on slide 3 with Principal box and three department boxes
Task ID: impress_teach_084
Domain: libreoffice_impress
Scoring:
  Component 1: Principal rectangle with correct fill (#1565C0) and white text (0.25)
  Component 2: Math Dept. rectangle with correct fill (#2E7D32) and white text (0.20)
  Component 3: Science Dept. rectangle with correct fill (#E65100) and white text (0.20)
  Component 4: English Dept. rectangle with correct fill (#6A1B9A) and white text (0.20)
  Component 5: Connecting lines between Principal and departments (0.15)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_084'


def get_shape_text(shape):
    """Extract all text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    return "".join(para.text for para in shape.text_frame.paragraphs).strip()


def get_shape_fill_rgb(shape):
    """Get solid fill color as uppercase hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_text_color_rgb(shape):
    """Get the font color of the first non-empty run as uppercase hex string, or None."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                try:
                    if run.font.color.type is not None:
                        return str(run.font.color.rgb).upper()
                except Exception:
                    pass
    return None


def find_rectangles_on_slide(slide):
    """Find all AUTO_SHAPE rectangles on a slide, returning list of dicts with text, fill, text_color."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = get_shape_text(shape)
            if text:
                rects.append({
                    'text': text,
                    'fill': get_shape_fill_rgb(shape),
                    'text_color': get_text_color_rgb(shape),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                })
    return rects


def count_lines_on_slide(slide):
    """Count LINE/CONNECTOR shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        # LINE is shape_type 9, freeform connectors may also appear
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            count += 1
        # Also check for FREEFORM (5) connectors sometimes used
        try:
            if 'Connector' in shape.name or 'Line' in shape.name:
                if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                    count = max(count, count)  # already counted if LINE
        except Exception:
            pass
    return count


def verify_task(file_path):
    """
    Verify organizational chart on slide 3 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed, slide 3

    # Get all rectangles and lines on slide 3
    rects = find_rectangles_on_slide(slide3)
    line_count = count_lines_on_slide(slide3)

    print(f"INFO: Found {len(rects)} labeled rectangles on slide 3")
    print(f"INFO: Found {line_count} line/connector shapes on slide 3")
    for r in rects:
        print(f"  - text='{r['text']}', fill={r['fill']}, text_color={r['text_color']}")

    # Build lookup by text (case-insensitive, stripped)
    rect_by_text = {}
    for r in rects:
        key = r['text'].strip().lower()
        rect_by_text[key] = r

    # Component 1: Principal rectangle with fill #1565C0 and white text (0.25 points)
    try:
        principal = rect_by_text.get('principal')
        if principal is not None:
            fill_ok = principal['fill'] == '1565C0'
            text_color_ok = principal['text_color'] == 'FFFFFF'
            if fill_ok and text_color_ok:
                print(f"PASS: Component 1 — Principal box: fill={principal['fill']}, text_color={principal['text_color']} (0.25 pts)")
                total_score += 0.25
            elif fill_ok or text_color_ok:
                print(f"PARTIAL: Component 1 — Principal box: fill={principal['fill']} (expect 1565C0), text_color={principal['text_color']} (expect FFFFFF) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 1 — Principal box: fill={principal['fill']} (expect 1565C0), text_color={principal['text_color']} (expect FFFFFF)")
        else:
            print(f"FAIL: Component 1 — No 'Principal' rectangle found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Math Dept. rectangle with fill #2E7D32 and white text (0.20 points)
    try:
        math_dept = rect_by_text.get('math dept.') or rect_by_text.get('math dept')
        if math_dept is not None:
            fill_ok = math_dept['fill'] == '2E7D32'
            text_color_ok = math_dept['text_color'] == 'FFFFFF'
            if fill_ok and text_color_ok:
                print(f"PASS: Component 2 — Math Dept. box: fill={math_dept['fill']}, text_color={math_dept['text_color']} (0.20 pts)")
                total_score += 0.20
            elif fill_ok or text_color_ok:
                print(f"PARTIAL: Component 2 — Math Dept. box: fill={math_dept['fill']} (expect 2E7D32), text_color={math_dept['text_color']} (expect FFFFFF) (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 2 — Math Dept. box: fill={math_dept['fill']} (expect 2E7D32), text_color={math_dept['text_color']} (expect FFFFFF)")
        else:
            print(f"FAIL: Component 2 — No 'Math Dept.' rectangle found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Science Dept. rectangle with fill #E65100 and white text (0.20 points)
    try:
        sci_dept = rect_by_text.get('science dept.') or rect_by_text.get('science dept')
        if sci_dept is not None:
            fill_ok = sci_dept['fill'] == 'E65100'
            text_color_ok = sci_dept['text_color'] == 'FFFFFF'
            if fill_ok and text_color_ok:
                print(f"PASS: Component 3 — Science Dept. box: fill={sci_dept['fill']}, text_color={sci_dept['text_color']} (0.20 pts)")
                total_score += 0.20
            elif fill_ok or text_color_ok:
                print(f"PARTIAL: Component 3 — Science Dept. box: fill={sci_dept['fill']} (expect E65100), text_color={sci_dept['text_color']} (expect FFFFFF) (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 3 — Science Dept. box: fill={sci_dept['fill']} (expect E65100), text_color={sci_dept['text_color']} (expect FFFFFF)")
        else:
            print(f"FAIL: Component 3 — No 'Science Dept.' rectangle found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: English Dept. rectangle with fill #6A1B9A and white text (0.20 points)
    try:
        eng_dept = rect_by_text.get('english dept.') or rect_by_text.get('english dept')
        if eng_dept is not None:
            fill_ok = eng_dept['fill'] == '6A1B9A'
            text_color_ok = eng_dept['text_color'] == 'FFFFFF'
            if fill_ok and text_color_ok:
                print(f"PASS: Component 4 — English Dept. box: fill={eng_dept['fill']}, text_color={eng_dept['text_color']} (0.20 pts)")
                total_score += 0.20
            elif fill_ok or text_color_ok:
                print(f"PARTIAL: Component 4 — English Dept. box: fill={eng_dept['fill']} (expect 6A1B9A), text_color={eng_dept['text_color']} (expect FFFFFF) (0.08 pts)")
                total_score += 0.08
            else:
                print(f"FAIL: Component 4 — English Dept. box: fill={eng_dept['fill']} (expect 6A1B9A), text_color={eng_dept['text_color']} (expect FFFFFF)")
        else:
            print(f"FAIL: Component 4 — No 'English Dept.' rectangle found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: At least 3 connecting lines between Principal and departments (0.15 points)
    try:
        if line_count >= 3:
            print(f"PASS: Component 5 — Found {line_count} connector lines (need >= 3) (0.15 pts)")
            total_score += 0.15
        elif line_count >= 1:
            partial = round(0.15 * line_count / 3, 2)
            print(f"PARTIAL: Component 5 — Found {line_count} connector lines (need >= 3) ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No connector lines found on slide 3")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
