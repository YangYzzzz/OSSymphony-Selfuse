"""
Initial Setup: Build a 5-slide school presentation with slide 3 empty (title only)
Task ID: impress_teach_084
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_084'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_text(slide, text):
    """Set the title placeholder text on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_body_text(slide, lines):
    """Add text lines to the body placeholder (index 1) if it exists."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Riverside Academy"
    slide1.placeholders[1].text = "Strategic Overview 2025-2026"

    # --- Slide 2: School Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide2, "School Overview")
    add_body_text(slide2, [
        "Founded in 1987 with a mission to inspire lifelong learning",
        "Serving 1,450 students across grades 6-12",
        "85 full-time faculty members across 8 departments",
        "Accredited by the Western Association of Schools",
        "Average class size: 22 students",
        "98.3% graduation rate (2024-2025)",
    ])

    # --- Slide 3: School Organization (EMPTY - title only) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title text box at the top
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "School Organization"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Academic Programs ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide4, "Academic Programs")
    add_body_text(slide4, [
        "Advanced Placement: 14 courses offered",
        "STEM Initiative: Robotics, coding, and engineering labs",
        "Arts Program: Visual arts, orchestra, theater, and choir",
        "Language Studies: Spanish, Mandarin, French, and Latin",
        "College readiness counseling from grade 9",
    ])

    # --- Slide 5: Contact & Enrollment ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_text(slide5, "Contact & Enrollment")
    add_body_text(slide5, [
        "Address: 2400 Riverside Drive, Maplewood, CA 94301",
        "Phone: (650) 555-0178",
        "Email: admissions@riversideacademy.edu",
        "Open House: March 15, 2026 | 9:00 AM - 1:00 PM",
        "Application deadline: April 30, 2026",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
