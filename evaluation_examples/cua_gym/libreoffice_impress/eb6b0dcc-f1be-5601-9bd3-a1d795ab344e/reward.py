"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 45 I want the title to show up in the exact “Dark Red 2” swatch (hex #B1001C), but with the shadow effect completely switched off. What steps do I follow in LibreOffice Impress to make that change?
Generated: 2025-09-10 12:49:32
Status: success
Model: azure-o3
Total Steps: 12
"""

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from lxml import etree
import zipfile
import os


def rgb_to_hex(rgb_color):
    """Convert python-pptx RGBColor -> HEX string (e.g. 'B1001C')."""
    try:
        return "%02X%02X%02X" % (rgb_color[0], rgb_color[1], rgb_color[2])
    except Exception:
        return None


def verify_task(file_path):
    """
    Reward script for the task:
      On slide 45, the title must be formatted in the exact colour
      ‘Dark Red 2’ (#B1001C) and have ALL shadow effects disabled.

    Scoring (progressive):
      • Correct colour   – 0.6
      • No shadow effect – 0.4
      Max = 1.0
    """
    max_score = 1.0
    score = 0.0
    colour_weight = 0.6
    shadow_weight = 0.4
    expected_hex = "B1001C"

    print(f"Checking file: {file_path}\n")

    # ---------- Prerequisite checks (no points) ----------
    if not os.path.exists(file_path):
        print("✗ File not found")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to open presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_index = 44  # slide 45 (0-based index)
    if len(prs.slides) <= slide_index:
        print(f"✗ Slide 45 missing (only {len(prs.slides)} slides)")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]

    # ---------- 1. Colour verification ----------
    title_shape = None
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            title_shape = sh
            break
    if title_shape is None:
        # Fallback: first non-empty text shape
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title_shape = sh
                break

    if title_shape is None:
        print("✗ Title shape not found on slide 45")
    else:
        all_runs_correct = True
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                colour = run.font.color
                if colour is None or colour.rgb is None:
                    all_runs_correct = False
                    print("   • Run without explicit RGB colour → incorrect")
                else:
                    actual_hex = rgb_to_hex(colour.rgb)
                    if actual_hex != expected_hex:
                        all_runs_correct = False
                        print(f"   • Colour mismatch: {actual_hex} (expected {expected_hex})")
        if all_runs_correct:
            score += colour_weight
            print(f"✓ Title colour is Dark Red 2 #{expected_hex} [{colour_weight}]")
        else:
            print("✗ Title colour incorrect")

    # ---------- 2. Shadow verification ----------
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            slide_name = f"ppt/slides/slide{slide_index + 1}.xml"  # archive is 1-based
            slide_xml = zf.read(slide_name)
            root = etree.fromstring(slide_xml)
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            shadow_elems = root.xpath('.//a:outerShdw | .//a:innerShdw', namespaces=ns)
            if not shadow_elems:
                score += shadow_weight
                print(f"✓ No shadow effects detected [{shadow_weight}]")
            else:
                print(f"✗ Found {len(shadow_elems)} shadow element(s) → shadow not disabled")
    except Exception as e:
        print(f"✗ Error inspecting slide XML: {e}")

    # ---------- Final ----------
    final_score = min(score, max_score)
    print(f"\nTotal score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == "__main__":
    verify_task("/home/user/on_slide_45_i_want_the_title_to_show_up_in_the_exact_dark_red_2_swatch_hex_b1001c_but_with_the_shado_golden.pptx")
