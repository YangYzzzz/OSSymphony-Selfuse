"""
FINAL REWARD SCRIPT - SUCCESS
Task: I've found some slides in my deck titled 'Backup Slides' that aren't needed anymore. Could you guide me on how I can remove those?
Generated: 2025-08-07 08:36:53
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_task(file_path):
    """
    Verifies that the PPTX file exists, loads it, and checks that no slides are titled 'Backup Slides'.
    Returns a progressive score up to 1.0 based on:
     - File existence (0.2)
     - Successful loading (0.3)
     - Absence of 'Backup Slides' titles (0.5)
    """
    print("Checking task completion...")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2 points)
    try:
        if os.path.exists(file_path):
            print(f"✓ File exists: {file_path} (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path} (0 points)")
    except Exception as e:
        print(f"✗ Error checking file existence: {e}")

    # Requirement 2: Load presentation (0.3 points)
    presentation = None
    try:
        presentation = Presentation(file_path)
        slide_count = len(presentation.slides)
        print(f"✓ Presentation loaded with {slide_count} slides (0.3 points)")
        total_score += 0.3
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        final = min(total_score, max_score)
        print(f"Final score: {final}")
        print(f"REWARD: {final}")
        return final

    # Requirement 3: No slides titled 'Backup Slides' (0.5 points)
    try:
        backup_count = 0
        for idx, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            if title_shape and title_shape.text:
                title_text = title_shape.text.strip().lower()
                if 'backup slide' in title_text:
                    backup_count += 1
                    print(f"✗ Found 'Backup Slides' title on slide {idx}: '{title_shape.text.strip()}'")
            else:
                # Also check all text shapes for occurrences
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        if 'backup slide' in shape.text.strip().lower():
                            backup_count += 1
                            print(f"✗ Found 'Backup Slides' text on slide {idx}: '{shape.text.strip()}'")
        if backup_count == 0:
            print("✓ No slides titled 'Backup Slides' found (0.5 points)")
            total_score += 0.5
        else:
            print(f"✗ {backup_count} slides still contain 'Backup Slides' (0 points for title removal)")
    except Exception as e:
        print(f"✗ Error checking slide titles: {e}")

    # Final score calculation
    final_score = min(total_score, max_score)
    print(f"Total score breakdown: {total_score}/{max_score}")
    print(f"Final score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# Execute verification
if __name__ == '__main__':
    file_path = '/home/user/ive_found_some_slides_in_my_deck_titled_backup_slides_that_arent_needed_anymore_could_you_guide_me_o.pptx'
    verify_task(file_path)
