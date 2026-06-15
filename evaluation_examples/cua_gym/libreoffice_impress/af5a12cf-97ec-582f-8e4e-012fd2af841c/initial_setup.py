"""
Initial Setup: Environmental Science presentation with inconsistent formatting
Task ID: impress_stu_026
Domain: libreoffice_impress

Creates a 10-slide presentation simulating merged content from different group
members, each with different fonts, sizes, colors, backgrounds, and bullet styles.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_026'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# --- Inconsistent formatting profiles for each "group member" ---
MEMBER_STYLES = [
    {  # Member A - dark bg, white title, large
        'bg': RGBColor(0x2C, 0x3E, 0x50),
        'title_font': 'Arial Black', 'title_size': Pt(36), 'title_color': RGBColor(0xFF, 0xFF, 0xFF), 'title_bold': True,
        'body_font': 'Arial', 'body_size': Pt(16), 'body_color': RGBColor(0xEC, 0xF0, 0xF1),
        'bullet_char': '-',
    },
    {  # Member B - white bg, blue title, serif body
        'bg': RGBColor(0xFF, 0xFF, 0xFF),
        'title_font': 'Times New Roman', 'title_size': Pt(28), 'title_color': RGBColor(0x00, 0x70, 0xC0), 'title_bold': False,
        'body_font': 'Times New Roman', 'body_size': Pt(14), 'body_color': RGBColor(0x00, 0x00, 0x00),
        'bullet_char': '•',
    },
    {  # Member C - green tint bg, dark green title
        'bg': RGBColor(0xE8, 0xF5, 0xE9),
        'title_font': 'Liberation Sans', 'title_size': Pt(30), 'title_color': RGBColor(0x1B, 0x5E, 0x20), 'title_bold': True,
        'body_font': 'Liberation Sans', 'body_size': Pt(20), 'body_color': RGBColor(0x33, 0x33, 0x33),
        'bullet_char': '>>',
    },
    {  # Member D - light blue bg, orange title
        'bg': RGBColor(0xE3, 0xF2, 0xFD),
        'title_font': 'DejaVu Sans', 'title_size': Pt(34), 'title_color': RGBColor(0xE6, 0x51, 0x00), 'title_bold': True,
        'body_font': 'DejaVu Sans', 'body_size': Pt(17), 'body_color': RGBColor(0x21, 0x21, 0x21),
        'bullet_char': '○',
    },
    {  # Member E - cream bg, brown title, italic
        'bg': RGBColor(0xFD, 0xF2, 0xE9),
        'title_font': 'Liberation Serif', 'title_size': Pt(26), 'title_color': RGBColor(0x79, 0x55, 0x48), 'title_bold': False,
        'body_font': 'Liberation Serif', 'body_size': Pt(15), 'body_color': RGBColor(0x4E, 0x34, 0x2E),
        'bullet_char': '→',
    },
]

# Slide content for Environmental Science presentation
SLIDES = [
    {
        'title': 'Environmental Science: Our Planet in Focus',
        'body': [
            'A comprehensive study of ecosystems, climate, and sustainability',
            'Prepared by: Group 7 - Fall 2025',
            'Professor: Dr. Elena Whitfield',
        ],
        'member': 0,
    },
    {
        'title': 'Climate Change Overview',
        'body': [
            'Global temperatures have risen 1.1 degrees C since pre-industrial era',
            'CO2 levels exceeded 420 ppm in 2024',
            'Arctic sea ice declining at 13% per decade',
            'Sea levels rising approximately 3.3 mm per year',
        ],
        'member': 1,
    },
    {
        'title': 'Biodiversity and Ecosystem Health',
        'body': [
            'Over 1 million species face extinction risk',
            'Coral reef coverage declined 50% since 1950',
            'Tropical deforestation: 4.7 million hectares annually',
            'Pollinator populations declining in 75% of crop types',
        ],
        'member': 2,
    },
    {
        'title': 'Water Resources and Management',
        'body': [
            '2.2 billion people lack safe drinking water access',
            'Agriculture accounts for 70% of global freshwater use',
            'Groundwater depletion accelerating in major aquifers',
            'Desalination capacity growing at 8% per year globally',
        ],
        'member': 3,
    },
    {
        'title': 'Renewable Energy Transition',
        'body': [
            'Solar PV costs dropped 89% between 2010 and 2023',
            'Wind energy now cheapest new electricity source in many regions',
            'Global renewable capacity reached 3,372 GW in 2023',
            'Battery storage deployments doubled year over year',
        ],
        'member': 4,
    },
    {
        'title': 'Air Quality and Pollution',
        'body': [
            '99% of the global population breathes air exceeding WHO limits',
            'Particulate matter PM2.5 causes 4.2 million premature deaths yearly',
            'Nitrogen dioxide levels improving in urban centers with EV adoption',
            'Indoor air pollution remains leading health risk in developing nations',
        ],
        'member': 0,
    },
    {
        'title': 'Sustainable Agriculture Practices',
        'body': [
            'Organic farming area expanded to 76.4 million hectares globally',
            'Precision agriculture reduces fertilizer use by 15-20%',
            'Vertical farming yields 300x more per square meter',
            'Agroforestry sequesters 2-4 tonnes CO2 per hectare per year',
        ],
        'member': 1,
    },
    {
        'title': 'Ocean Conservation Challenges',
        'body': [
            'Ocean absorbs 30% of CO2 emissions, increasing acidity',
            'Plastic pollution: 11 million tonnes enter oceans annually',
            'Marine protected areas cover only 8.2% of ocean surface',
            'Deep-sea mining threatens unexplored ecosystems',
        ],
        'member': 2,
    },
    {
        'title': 'Urban Environmental Planning',
        'body': [
            '68% of world population will live in cities by 2050',
            'Green infrastructure reduces urban heat island by 2-3 degrees C',
            'Public transit expansion linked to 45% emission reduction per capita',
            'Building sector accounts for 39% of energy-related CO2 emissions',
        ],
        'member': 3,
    },
    {
        'title': 'Conclusions and Future Outlook',
        'body': [
            'Integrated approaches essential for addressing climate change',
            'Technology and policy must work together for net-zero by 2050',
            'Community engagement drives successful conservation outcomes',
            'Continued research funding critical for evidence-based solutions',
        ],
        'member': 4,
    },
]


def set_bullet(paragraph, char):
    """Set a custom bullet character on a paragraph via XML."""
    pPr = paragraph._p.get_or_add_pPr()
    # Remove any existing bullet elements
    for tag in ['buChar', 'buAutoNum', 'buNone']:
        existing = pPr.find(qn(f'a:{tag}'))
        if existing is not None:
            pPr.remove(existing)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
    pPr.append(buChar)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_info in SLIDES:
        style = MEMBER_STYLES[slide_info['member']]
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = style['bg']

        # Title textbox
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info['title']
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = style['title_font']
        run.font.size = style['title_size']
        run.font.color.rgb = style['title_color']
        run.font.bold = style['title_bold']

        # Body textbox with bullet points
        bodyBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.5))
        btf = bodyBox.text_frame
        btf.word_wrap = True

        for i, line in enumerate(slide_info['body']):
            if i == 0:
                para = btf.paragraphs[0]
            else:
                para = btf.add_paragraph()
            para.text = line
            para.space_after = Pt(8)
            set_bullet(para, style['bullet_char'])
            run = para.runs[0]
            run.font.name = style['body_font']
            run.font.size = style['body_size']
            run.font.color.rgb = style['body_color']

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
