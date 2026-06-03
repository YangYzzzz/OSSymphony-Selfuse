"""
Reward Script: Apply consistent color scheme to presentation
Task ID: impress_stu_026
Domain: libreoffice_impress
Scoring:
  Component 1: Slide backgrounds all #FAFAFA (0.25)
  Component 2: Title text formatting - #2C3E50, 32pt, bold (0.30)
  Component 3: Body text formatting - #34495E, 18pt (0.25)
  Component 4: Square bullet characters on all body text (0.20)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_026'

# Expected values from task
EXPECTED_BG_COLOR = 'FAFAFA'
EXPECTED_TITLE_COLOR = '2C3E50'
EXPECTED_TITLE_SIZE = 406400   # 32pt in EMU
EXPECTED_BODY_COLOR = '34495E'
EXPECTED_BODY_SIZE = 228600    # 18pt in EMU
EXPECTED_BULLET_CHAR = '\u25A0'  # Black square


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides == 0:
        print("CRITICAL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: All slide backgrounds are #FAFAFA (0.25 points)
    try:
        bg_pass = 0
        bg_total = num_slides
        for i, slide in enumerate(prs.slides):
            fill = slide.background.fill
            if fill.type == 1:  # SOLID fill
                color_str = str(fill.fore_color.rgb)
                if color_str == EXPECTED_BG_COLOR:
                    bg_pass += 1
                else:
                    print(f"  FAIL: Slide {i+1} background is {color_str}, expected {EXPECTED_BG_COLOR}")
            else:
                print(f"  FAIL: Slide {i+1} background is not solid fill (type={fill.type})")

        if bg_pass == bg_total:
            print(f"PASS: Component 1 - All {bg_total} slides have background #{EXPECTED_BG_COLOR} (0.25 pts)")
            total_score += 0.25
        elif bg_pass > 0:
            partial = 0.25 * (bg_pass / bg_total)
            print(f"PARTIAL: Component 1 - {bg_pass}/{bg_total} slides have correct background ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No slides have background #{EXPECTED_BG_COLOR}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Title text is #2C3E50, 32pt (406400 EMU), bold (0.30 points)
    # Titles are identified as the first text shape with content on each slide (typically "TextBox 2")
    try:
        title_pass = 0
        title_total = 0
        for i, slide in enumerate(prs.slides):
            # Find title shape: look for shapes with text, pick the one that appears to be a title
            # (shorter text, typically the first text-bearing shape after placeholder)
            title_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        title_shape = shape
                        break  # first text shape with content is the title

            if title_shape is None:
                continue

            title_total += 1
            slide_ok = True
            for para in title_shape.text_frame.paragraphs:
                runs = [r for r in para.runs if (r.text or "").strip()]
                for run in runs:
                    # Check color
                    try:
                        if run.font.color.type is not None:
                            color_str = str(run.font.color.rgb)
                        else:
                            color_str = 'None'
                    except:
                        color_str = 'error'
                    if color_str != EXPECTED_TITLE_COLOR:
                        slide_ok = False
                        print(f"  FAIL: Slide {i+1} title run color={color_str}, expected {EXPECTED_TITLE_COLOR}")

                    # Check size
                    if run.font.size != EXPECTED_TITLE_SIZE:
                        slide_ok = False
                        print(f"  FAIL: Slide {i+1} title run size={run.font.size}, expected {EXPECTED_TITLE_SIZE}")

                    # Check bold
                    is_bold = run.font.bold if run.font.bold is not None else False
                    if not is_bold:
                        slide_ok = False
                        print(f"  FAIL: Slide {i+1} title run bold={run.font.bold}, expected True")

            if slide_ok:
                title_pass += 1

        if title_total > 0 and title_pass == title_total:
            print(f"PASS: Component 2 - All {title_total} titles are #2C3E50, 32pt, bold (0.30 pts)")
            total_score += 0.30
        elif title_pass > 0:
            partial = 0.30 * (title_pass / title_total)
            print(f"PARTIAL: Component 2 - {title_pass}/{title_total} titles correctly formatted ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 - No titles have correct formatting")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Body text is #34495E, 18pt (228600 EMU) (0.25 points)
    try:
        body_pass = 0
        body_total = 0
        for i, slide in enumerate(prs.slides):
            # Body text is in shapes after the title shape
            text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
            if len(text_shapes) < 2:
                continue  # no body text shape

            # Body shapes are all text shapes except the first (title)
            for body_shape in text_shapes[1:]:
                for para in body_shape.text_frame.paragraphs:
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    for run in runs:
                        body_total += 1
                        run_ok = True

                        # Check color
                        try:
                            if run.font.color.type is not None:
                                color_str = str(run.font.color.rgb)
                            else:
                                color_str = 'None'
                        except:
                            color_str = 'error'
                        if color_str != EXPECTED_BODY_COLOR:
                            run_ok = False

                        # Check size
                        if run.font.size != EXPECTED_BODY_SIZE:
                            run_ok = False

                        if run_ok:
                            body_pass += 1
                        else:
                            if body_total <= 5:  # limit verbose output
                                print(f"  FAIL: Slide {i+1} body run color={color_str} size={run.font.size}")

        if body_total > 0 and body_pass == body_total:
            print(f"PASS: Component 3 - All {body_total} body text runs are #34495E 18pt (0.25 pts)")
            total_score += 0.25
        elif body_pass > 0:
            partial = 0.25 * (body_pass / body_total)
            print(f"PARTIAL: Component 3 - {body_pass}/{body_total} body runs correctly formatted ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - No body text runs have correct formatting")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Square bullet characters on all bullet points (0.20 points)
    # Check via XML for buChar elements
    try:
        ns_a = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        bullet_pass = 0
        bullet_total = 0

        with zipfile.ZipFile(file_path, 'r') as zf:
            for si in range(1, num_slides + 1):
                fname = f'ppt/slides/slide{si}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        for para in root.findall('.//a:p', ns_a):
                            # Only check paragraphs with text content
                            text = "".join(t.text or "" for t in para.findall('.//a:t', ns_a))
                            if not text.strip():
                                continue

                            pPr = para.find('a:pPr', ns_a)
                            if pPr is None:
                                continue

                            buChar = pPr.find('a:buChar', ns_a)
                            if buChar is not None:
                                bullet_total += 1
                                char = buChar.get('char', '')
                                # Square bullet: U+25A0 (black square) or U+25AA (small black square)
                                if char in ('\u25A0', '\u25AA'):
                                    bullet_pass += 1
                                else:
                                    if bullet_total <= 5:
                                        print(f"  FAIL: Slide {si} has bullet char '{char}' (U+{ord(char):04X}), expected square")
                except KeyError:
                    pass

        if bullet_total > 0 and bullet_pass == bullet_total:
            print(f"PASS: Component 4 - All {bullet_total} bullet points use square bullets (0.20 pts)")
            total_score += 0.20
        elif bullet_pass > 0:
            partial = 0.20 * (bullet_pass / bullet_total)
            print(f"PARTIAL: Component 4 - {bullet_pass}/{bullet_total} bullets are square ({partial:.3f} pts)")
            total_score += partial
        elif bullet_total == 0:
            print(f"FAIL: Component 4 - No bullet characters found in presentation")
        else:
            print(f"FAIL: Component 4 - No square bullets found (0/{bullet_total})")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
