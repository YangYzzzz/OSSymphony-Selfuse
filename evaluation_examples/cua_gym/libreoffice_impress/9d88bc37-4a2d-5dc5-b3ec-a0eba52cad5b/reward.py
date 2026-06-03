"""
Reward Script: Change the background of every single slide to solid color #2C3E50
Task ID: osworld_impress_all_slides_background_009
Domain: libreoffice_impress
Scoring:
  Component 1: Number of slides changed to #2C3E50 solid background (0.7 pts - proportional)
  Component 2: All 12 slides correctly changed to #2C3E50 solid background (0.3 pts bonus for full completion)
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'  # VM path — all reward scripts run on the VM
TASK_ID = 'osworld_impress_all_slides_background_009'

TARGET_COLOR = '2C3E50'   # Expected background color (hex, uppercase, no '#')
EXPECTED_SLIDE_COUNT = 12  # Ground truth from task context


def get_slide_background_rgb(slide):
    """
    Returns the RGBColor of a slide's background fill.
    Handles solid fill (type 1) and inherited from master (type 5).
    Returns None if no color can be determined.
    """
    fill = slide.background.fill
    if fill.type == 1:  # MSO_FILL.SOLID
        return fill.fore_color.rgb
    elif fill.type == 5:  # MSO_FILL.BACKGROUND (inherited from master)
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return master_fill.fore_color.rgb
        except Exception:
            pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0

    Scoring breakdown:
    - Component 1: Proportion of slides with #2C3E50 solid background (0.7 points, proportional)
    - Component 2: All EXPECTED_SLIDE_COUNT slides are correctly changed (0.3 points bonus)
    Total: 1.0 when all slides correctly set to #2C3E50
    """
    total_score = 0.0

    # Load the presentation file (precondition gate)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: check that the slide count is as expected
    num_slides = len(prs.slides)
    print(f"INFO: Presentation has {num_slides} slides (expected: {EXPECTED_SLIDE_COUNT})")

    if num_slides == 0:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    # --- Component 1: Proportion of slides with #2C3E50 solid background (0.7 points) ---
    # Check each slide for the correct background color
    slides_with_correct_bg = 0
    slides_checked = min(num_slides, EXPECTED_SLIDE_COUNT)

    for i in range(slides_checked):
        slide = prs.slides[i]
        try:
            rgb = get_slide_background_rgb(slide)
            if rgb is not None and str(rgb).upper() == TARGET_COLOR:
                slides_with_correct_bg += 1
            else:
                actual_color = str(rgb).upper() if rgb is not None else "None/Inherited"
                print(f"FAIL: Slide {i+1} background is '{actual_color}', expected '{TARGET_COLOR}'")
        except Exception as e:
            print(f"ERROR: Could not check background for slide {i+1}: {e}")

    proportion_correct = slides_with_correct_bg / slides_checked
    component1_score = round(proportion_correct * 0.7, 4)

    if slides_with_correct_bg == slides_checked:
        print(f"PASS: Component 1 — All {slides_checked} slides have #2C3E50 solid background ({component1_score:.4f} pts)")
        total_score += component1_score
    elif slides_with_correct_bg > 0:
        print(f"PARTIAL: Component 1 — {slides_with_correct_bg}/{slides_checked} slides have correct background "
              f"({component1_score:.4f} pts)")
        total_score += component1_score
    else:
        print(f"FAIL: Component 1 — 0/{slides_checked} slides have correct background (0.0 pts)")

    # --- Component 2: All slides correctly changed AND slide count matches (0.3 points) ---
    # This gives bonus points only when every single slide is correctly set
    try:
        all_slides_correct = (slides_with_correct_bg == EXPECTED_SLIDE_COUNT
                              and num_slides == EXPECTED_SLIDE_COUNT)

        if all_slides_correct:
            print(f"PASS: Component 2 — All {EXPECTED_SLIDE_COUNT} slides correctly changed to #2C3E50 (0.3 pts)")
            total_score += 0.3
        else:
            if num_slides != EXPECTED_SLIDE_COUNT:
                print(f"FAIL: Component 2 — Expected {EXPECTED_SLIDE_COUNT} slides, found {num_slides}")
            else:
                print(f"FAIL: Component 2 — Only {slides_with_correct_bg}/{EXPECTED_SLIDE_COUNT} slides "
                      f"have the correct background color (need all {EXPECTED_SLIDE_COUNT})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.4f}/1.0")
    print(f"REWARD: {final_score:.1f}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
