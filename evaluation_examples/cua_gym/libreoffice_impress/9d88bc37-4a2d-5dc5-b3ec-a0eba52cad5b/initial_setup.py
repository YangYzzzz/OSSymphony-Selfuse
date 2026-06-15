"""
Initial Setup: 12-slide annual report deck with white backgrounds
Task ID: osworld_impress_all_slides_background_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_all_slides_background_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_white_background(slide):
    """Set solid white background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_text(slide, title_text, subtitle_text=None, title_size=36, subtitle_size=18):
    """Add title and optional subtitle text boxes to a slide."""
    prs_width = Inches(10)
    prs_height = Inches(7.5)

    # Title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)  # dark blue-gray text for contrast

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.runs[0]
        run2.font.size = Pt(subtitle_size)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_content_box(slide, title_text, body_lines, top=Inches(1.5)):
    """Add a heading + body text content block to a slide."""
    # Section heading
    txHead = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.7))
    tf = txHead.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x5F, 0x7A)

    # Body text box
    txBody = slide.shapes.add_textbox(Inches(0.5), top + Inches(0.75), Inches(9), Inches(4.5))
    tf2 = txBody.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        if i == 0:
            para = tf2.paragraphs[0]
        else:
            para = tf2.add_paragraph()
        para.text = line
        para.level = 0
        run = para.runs[0] if para.runs else para.add_run()
        if not para.runs:
            run = para.add_run()
            run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(blank_layout)
    set_white_background(slide1)
    add_title_text(slide1,
                   "Nexus Dynamics Inc.",
                   subtitle_text="Annual Report 2024 — Powering Growth Through Innovation",
                   title_size=40, subtitle_size=20)

    # ---- Slide 2: Table of Contents ----
    slide2 = prs.slides.add_slide(blank_layout)
    set_white_background(slide2)
    add_content_box(slide2, "Table of Contents", [
        "1. Executive Summary",
        "2. Financial Highlights",
        "3. Revenue Breakdown",
        "4. Product Portfolio",
        "5. Market Analysis",
        "6. Operational Excellence",
        "7. Customer Success",
        "8. Research & Development",
        "9. Global Expansion",
        "10. Sustainability Initiatives",
        "11. Leadership Team",
        "12. Outlook & Strategy 2025",
    ])

    # ---- Slide 3: Executive Summary ----
    slide3 = prs.slides.add_slide(blank_layout)
    set_white_background(slide3)
    add_content_box(slide3, "Executive Summary", [
        "Nexus Dynamics achieved record-breaking performance in fiscal year 2024, "
        "delivering $2.4B in total revenue — a 18.3% year-over-year increase.",
        "",
        "Key Achievements:",
        "  • Net income reached $412M, up 22% from 2023",
        "  • Expanded operations into 6 new international markets",
        "  • Launched 4 flagship product lines across enterprise and consumer segments",
        "  • Customer satisfaction score rose to 91% (from 87% in 2023)",
        "  • Workforce grew by 1,240 employees to a total of 9,800 globally",
    ])

    # ---- Slide 4: Financial Highlights ----
    slide4 = prs.slides.add_slide(blank_layout)
    set_white_background(slide4)
    add_content_box(slide4, "Financial Highlights", [
        "Revenue:            $2,400M   (+18.3% YoY)",
        "Gross Profit:       $1,056M   (+21.1% YoY)",
        "Operating Income:   $532M     (+19.6% YoY)",
        "Net Income:         $412M     (+22.0% YoY)",
        "EPS (Diluted):      $4.87     (+20.8% YoY)",
        "Cash & Equivalents: $890M",
        "Total Assets:       $5,120M",
        "Return on Equity:   17.4%",
    ])

    # ---- Slide 5: Revenue Breakdown ----
    slide5 = prs.slides.add_slide(blank_layout)
    set_white_background(slide5)
    add_content_box(slide5, "Revenue Breakdown by Segment", [
        "Enterprise Software Solutions:  $940M   (39.2%)",
        "Cloud & SaaS Services:          $720M   (30.0%)",
        "Hardware & IoT Devices:         $480M   (20.0%)",
        "Professional Services:          $260M   (10.8%)",
        "",
        "Top Revenue Regions:",
        "  North America:    $1,056M (44%)",
        "  Europe:           $720M   (30%)",
        "  Asia-Pacific:     $480M   (20%)",
        "  Rest of World:    $144M   (6%)",
    ])

    # ---- Slide 6: Product Portfolio ----
    slide6 = prs.slides.add_slide(blank_layout)
    set_white_background(slide6)
    add_content_box(slide6, "Product Portfolio 2024", [
        "NexusOS Enterprise 5.0  — Advanced OS platform for data centers (launched Q1 2024)",
        "CloudBridge Pro 3.2     — Hybrid cloud management suite (launched Q2 2024)",
        "IotSense Guardian       — Industrial IoT security module (launched Q3 2024)",
        "DataVault Analytics     — Real-time business intelligence platform (launched Q4 2024)",
        "",
        "Legacy Product Performance:",
        "  NexusOS Enterprise 4.x — 2,400 active enterprise clients",
        "  CloudBridge Pro 2.x    — 5,800 subscriptions renewed",
    ])

    # ---- Slide 7: Market Analysis ----
    slide7 = prs.slides.add_slide(blank_layout)
    set_white_background(slide7)
    add_content_box(slide7, "Market Analysis", [
        "Total Addressable Market (TAM): $85B by 2027 (CAGR 14.2%)",
        "Nexus Dynamics Market Share:    2.8% global (up from 2.4% in 2023)",
        "",
        "Competitive Landscape:",
        "  • Nexus Dynamics holds #3 position in enterprise software globally",
        "  • Key differentiators: AI integration, security, scalability",
        "  • Competitors: Apex Systems (21%), OrbitTech (18%), Nexus (13%), others (48%)",
        "",
        "Growth Drivers:",
        "  Digital transformation acceleration post-pandemic",
        "  Increasing demand for hybrid cloud solutions",
        "  Regulatory requirements driving security investments",
    ])

    # ---- Slide 8: Operational Excellence ----
    slide8 = prs.slides.add_slide(blank_layout)
    set_white_background(slide8)
    add_content_box(slide8, "Operational Excellence", [
        "Supply Chain Optimization: 15% reduction in logistics costs via AI-driven routing",
        "Manufacturing Efficiency: 8.2% improvement in output per employee",
        "Quality Metrics: 99.2% on-time delivery rate for hardware shipments",
        "",
        "Technology Infrastructure:",
        "  • 12 global data centers (3 new in 2024: Singapore, Frankfurt, São Paulo)",
        "  • 99.97% system uptime across all cloud services",
        "  • 42 petabytes of managed data",
        "",
        "Process Certifications:",
        "  ISO 27001, ISO 9001, SOC 2 Type II, PCI DSS Level 1",
    ])

    # ---- Slide 9: Customer Success ----
    slide9 = prs.slides.add_slide(blank_layout)
    set_white_background(slide9)
    add_content_box(slide9, "Customer Success", [
        "Total Customers: 24,600 (enterprise + SMB + consumer)",
        "Net Promoter Score (NPS): 67 (industry avg: 42)",
        "Customer Retention Rate: 93.4%",
        "New Enterprise Clients Acquired: 380",
        "",
        "Notable Client Wins in 2024:",
        "  • GlobalBank Financial Services — 3-year cloud transformation contract ($85M)",
        "  • Meridian Healthcare Group — NexusOS Enterprise rollout (18,000 seats)",
        "  • Vortex Logistics Corp — IotSense Guardian deployment (12,000 devices)",
        "  • Skyline Media Consortium — DataVault Analytics subscription (5-year)",
    ])

    # ---- Slide 10: Research & Development ----
    slide10 = prs.slides.add_slide(blank_layout)
    set_white_background(slide10)
    add_content_box(slide10, "Research & Development", [
        "R&D Investment: $340M (14.2% of revenue)",
        "Patents Filed in 2024: 127 new patents",
        "Total Patent Portfolio: 1,840 active patents",
        "R&D Headcount: 2,200 engineers and researchers",
        "",
        "Key Research Areas:",
        "  • Artificial Intelligence & Machine Learning integration",
        "  • Quantum computing readiness framework",
        "  • Zero-trust security architecture",
        "  • Edge computing and 5G integration",
        "",
        "Innovation Labs: Austin TX, Zurich, Tokyo, Bangalore",
    ])

    # ---- Slide 11: Leadership Team ----
    slide11 = prs.slides.add_slide(blank_layout)
    set_white_background(slide11)
    add_content_box(slide11, "Leadership Team", [
        "Victoria Hartwell   — Chief Executive Officer (CEO), 14 years at Nexus",
        "Raymond Osei        — Chief Financial Officer (CFO), joined 2019",
        "Mei-Lin Zhao        — Chief Technology Officer (CTO), PhD MIT CSAIL",
        "Derek Flannery      — Chief Operating Officer (COO), ex-McKinsey",
        "Priya Nambiar       — Chief Marketing Officer (CMO), former Google VP",
        "Samuel Kowalczyk    — SVP Engineering, 200+ team leads",
        "Isabelle Fontaine   — SVP Global Sales, $1.8B revenue owned",
        "Kwame Asantewaa     — SVP People & Culture, 9,800 employees globally",
    ])

    # ---- Slide 12: Outlook & Strategy 2025 ----
    slide12 = prs.slides.add_slide(blank_layout)
    set_white_background(slide12)
    add_content_box(slide12, "Outlook & Strategy 2025", [
        "Revenue Target: $2.85B (+18.75% from 2024)",
        "Net Income Target: $500M",
        "New Markets: Brazil, India, Middle East expansion",
        "",
        "Strategic Priorities:",
        "  1. Accelerate AI-native product line (NexusAI Suite — launch Q2 2025)",
        "  2. Complete acquisition of DataStream Analytics ($320M deal, pending)",
        "  3. Expand CloudBridge subscription base to 10,000+ by Q4 2025",
        "  4. Launch Partner Ecosystem Program targeting 500 certified ISVs",
        "  5. Achieve carbon neutrality across all owned data centers by Dec 2025",
        "",
        "Long-Term Vision: Become Top 2 enterprise software provider globally by 2028",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open the file in LibreOffice Impress with DISPLAY=:0
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
