"""
Initial Setup: Create a lecture presentation with 8 slides, no progress indicators
Task ID: impress_teach_080
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16,
                      font_name="Arial", color=None):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
    ACCENT_ORANGE = RGBColor(0xFF, 0xA7, 0x26)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "Introduction to Machine Learning", font_size=36, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                 "Module 3: Supervised Learning Fundamentals", font_size=22,
                 color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                 "Dr. Elena Rodriguez  |  Stanford CS229  |  Spring 2025", font_size=16,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide2.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Module Overview", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide2, Inches(0.8), Inches(1.5), Inches(11), Inches(5.0),
                      [
                          "What is supervised learning and how does it differ from unsupervised methods?",
                          "Key terminology: features, labels, training set, test set, validation",
                          "The hypothesis space and model selection criteria",
                          "Bias-variance tradeoff and its practical implications",
                          "Real-world applications across healthcare, finance, and autonomous systems",
                      ], font_size=18, color=LIGHT_GRAY)

    # --- Slide 3: Core Concepts ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide3.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Core Concepts of Supervised Learning", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide3, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
                      [
                          "Training data: labeled pairs (x, y) drawn from unknown distribution D",
                          "Learning algorithm: maps training set to a hypothesis h: X -> Y",
                          "Loss function: quantifies prediction error (e.g., MSE, cross-entropy)",
                          "Empirical risk minimization: find h that minimizes average loss",
                      ], font_size=16, color=LIGHT_GRAY)
    add_text_box(slide3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.0),
                 "Example:\nGiven 10,000 chest X-rays labeled normal/pneumonia,\n"
                 "train a classifier f(x) that predicts the diagnosis\n"
                 "for new, unseen patient images.\n\n"
                 "Training accuracy: 97.3%\nValidation accuracy: 94.1%",
                 font_size=14, color=ACCENT_ORANGE)

    # --- Slide 4: Regression Models ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide4.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Linear and Polynomial Regression", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide4, Inches(0.8), Inches(1.5), Inches(11), Inches(5.0),
                      [
                          "Linear regression: h(x) = w^T x + b, minimizes squared error",
                          "Closed-form solution: w* = (X^T X)^{-1} X^T y (Normal Equation)",
                          "Gradient descent alternative: iteratively update w -= alpha * grad(L)",
                          "Polynomial features: map x -> [1, x, x^2, ..., x^d] for nonlinear fits",
                          "Regularization (Ridge, Lasso) prevents overfitting on high-degree polynomials",
                          "Evaluation: R-squared, adjusted R-squared, residual analysis",
                      ], font_size=16, color=LIGHT_GRAY)

    # --- Slide 5: Classification ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide5.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Classification Algorithms", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide5, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0),
                      [
                          "Logistic regression: P(y=1|x) = sigma(w^T x + b)",
                          "Decision trees: recursive feature-space partitioning",
                          "Support Vector Machines: maximize margin between classes",
                          "k-Nearest Neighbors: classify by majority vote of neighbors",
                          "Naive Bayes: apply Bayes' theorem with feature independence",
                      ], font_size=16, color=LIGHT_GRAY)

    # Table on right side
    table_shape = slide5.shapes.add_table(5, 3, Inches(7.0), Inches(1.8), Inches(5.5), Inches(3.0))
    table = table_shape.table
    headers = ["Algorithm", "Accuracy", "Training Time"]
    data = [
        ["Logistic Reg.", "89.2%", "0.3 sec"],
        ["Decision Tree", "91.5%", "1.2 sec"],
        ["SVM (RBF)", "93.8%", "14.7 sec"],
        ["k-NN (k=5)", "88.6%", "0.01 sec"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 6: Model Evaluation ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide6.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Model Evaluation Techniques", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide6, Inches(0.8), Inches(1.5), Inches(11), Inches(5.0),
                      [
                          "Train/test split: typically 80/20 or 70/30 for initial evaluation",
                          "k-Fold cross-validation: partition data into k subsets, rotate test fold",
                          "Confusion matrix: true positives, false positives, true negatives, false negatives",
                          "Precision = TP / (TP + FP), Recall = TP / (TP + FN)",
                          "F1 score = 2 * (Precision * Recall) / (Precision + Recall)",
                          "ROC curve and AUC: plot true positive rate vs false positive rate",
                      ], font_size=16, color=LIGHT_GRAY)

    # --- Slide 7: Practical Tips ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide7.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "Practical Implementation Tips", font_size=30, bold=True, color=WHITE)
    add_bullet_points(slide7, Inches(0.8), Inches(1.5), Inches(11), Inches(5.0),
                      [
                          "Always start with exploratory data analysis (EDA) before modeling",
                          "Normalize or standardize features to comparable scales",
                          "Handle missing values: imputation, deletion, or indicator variables",
                          "Use stratified sampling for imbalanced classification tasks",
                          "Monitor learning curves to diagnose underfitting vs overfitting",
                          "Document your pipeline: data preprocessing, feature engineering, model selection",
                      ], font_size=16, color=LIGHT_GRAY)

    # --- Slide 8: Summary & Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide8.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    add_text_box(slide8, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
                 "Summary & Next Steps", font_size=32, bold=True,
                 color=WHITE, alignment=PP_ALIGN.CENTER)
    add_bullet_points(slide8, Inches(1.5), Inches(3.0), Inches(10), Inches(3.5),
                      [
                          "Supervised learning uses labeled data to learn predictive mappings",
                          "Regression predicts continuous values; classification predicts discrete labels",
                          "Model evaluation requires careful train/test methodology",
                          "Next module: Unsupervised Learning and Clustering (Module 4)",
                          "Assignment due: Implement linear regression on the housing dataset by Friday",
                      ], font_size=18, color=LIGHT_GRAY)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
