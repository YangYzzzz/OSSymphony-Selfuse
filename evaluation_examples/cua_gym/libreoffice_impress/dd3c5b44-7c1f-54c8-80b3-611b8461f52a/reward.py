"""
Reward Script: Progress indicator on slides 2-6
Task ID: impress_teach_080
Domain: libreoffice_impress
Scoring: 5 components (one per slide 2-6), each worth 0.2 points.
Each component checks for correct progress bar rectangles at the bottom.
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_080'

# Expected green fill percentages for slides 2-6
EXPECTED_GREEN_PCT = {
    2: 20,
    3: 40,
    4: 60,
    5: 80,
    6: 100,
}

GREEN_COLOR = '4CAF50'
GRAY_COLOR = 'E0E0E0'
BAR_HEIGHT_EMU = 137160  # 0.15 inches
HEIGHT_TOLERANCE = 0.05  # 5% relative tolerance
PCT_TOLERANCE = 3.0  # percentage points tolerance for width checks


def get_shape_fill_rgb(shape):
    """Get the fill color of a shape as hex string, or None."""
    try:
        fill = shape.fill
        if fill.type is not None:
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    num_slides = len(prs.slides)

    if num_slides < 6:
        print(f"FAIL: Presentation has only {num_slides} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    # Expected bar position: at the very bottom of the slide
    expected_bar_top = slide_height - BAR_HEIGHT_EMU

    for slide_num in range(2, 7):
        slide_idx = slide_num - 1
        slide = prs.slides[slide_idx]
        expected_green_pct = EXPECTED_GREEN_PCT[slide_num]
        points = 0.2

        # Component: Slide N progress bar (0.2 points)
        try:
            # Find all AUTO_SHAPE rectangles near the bottom of the slide
            bar_shapes = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    # Check if shape is at the bottom (within tolerance of expected position)
                    top_diff = abs(shape.top - expected_bar_top)
                    height_ok = abs(shape.height - BAR_HEIGHT_EMU) / BAR_HEIGHT_EMU <= HEIGHT_TOLERANCE
                    top_ok = top_diff / slide_height <= 0.02  # within 2% of slide height
                    if height_ok and top_ok:
                        color = get_shape_fill_rgb(shape)
                        bar_shapes.append({
                            'name': shape.name,
                            'left': shape.left,
                            'width': shape.width,
                            'color': color,
                        })

            if not bar_shapes:
                print(f"FAIL: Slide {slide_num} — no progress bar rectangles found at bottom")
                continue

            # Identify green and gray portions
            green_shapes = [s for s in bar_shapes if s['color'] == GREEN_COLOR]
            gray_shapes = [s for s in bar_shapes if s['color'] == GRAY_COLOR]

            # For slide 6 (100% green), there should be no gray portion
            if expected_green_pct == 100:
                if not green_shapes:
                    print(f"FAIL: Slide {slide_num} — no green (#4CAF50) bar found")
                    continue
                # Check green covers full width
                green_total_width = sum(s['width'] for s in green_shapes)
                green_pct = green_total_width / slide_width * 100
                if abs(green_pct - 100) <= PCT_TOLERANCE:
                    # Check green starts at left edge
                    green_starts_at_zero = any(s['left'] <= slide_width * 0.01 for s in green_shapes)
                    if green_starts_at_zero:
                        print(f"PASS: Slide {slide_num} — 100% green bar ({green_pct:.1f}% width) ({points} pts)")
                        total_score += points
                    else:
                        print(f"FAIL: Slide {slide_num} — green bar does not start at left edge")
                else:
                    print(f"FAIL: Slide {slide_num} — green bar covers {green_pct:.1f}%, expected ~100%")
            else:
                # Need both green and gray portions
                if not green_shapes:
                    print(f"FAIL: Slide {slide_num} — no green (#4CAF50) bar found")
                    continue
                if not gray_shapes:
                    print(f"FAIL: Slide {slide_num} — no gray (#E0E0E0) bar found")
                    continue

                green_width = sum(s['width'] for s in green_shapes)
                gray_width = sum(s['width'] for s in gray_shapes)
                green_pct = green_width / slide_width * 100
                gray_pct = gray_width / slide_width * 100

                # Check green percentage matches expected
                green_ok = abs(green_pct - expected_green_pct) <= PCT_TOLERANCE
                # Check total coverage is ~100%
                total_coverage = green_pct + gray_pct
                coverage_ok = abs(total_coverage - 100) <= PCT_TOLERANCE * 2

                # Check green starts from left
                green_starts_at_zero = any(s['left'] <= slide_width * 0.01 for s in green_shapes)

                if green_ok and coverage_ok and green_starts_at_zero:
                    print(f"PASS: Slide {slide_num} — green {green_pct:.1f}% (expected {expected_green_pct}%), gray {gray_pct:.1f}% ({points} pts)")
                    total_score += points
                else:
                    reasons = []
                    if not green_ok:
                        reasons.append(f"green is {green_pct:.1f}%, expected {expected_green_pct}%")
                    if not coverage_ok:
                        reasons.append(f"total coverage {total_coverage:.1f}%, expected ~100%")
                    if not green_starts_at_zero:
                        reasons.append("green bar doesn't start at left edge")
                    print(f"FAIL: Slide {slide_num} — {'; '.join(reasons)}")

        except Exception as e:
            print(f"ERROR: Slide {slide_num} — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Main execution
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
