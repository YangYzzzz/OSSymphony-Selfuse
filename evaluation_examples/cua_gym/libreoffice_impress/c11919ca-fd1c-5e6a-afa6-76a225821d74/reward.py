"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 56 is driving me nuts—LibreOffice Impress keeps underlining every single word in the title because it’s set to the wrong language. How do I change that title’s language setting specifically to English (US) so the spell-checker stops screaming at me?
Generated: 2025-09-10 22:48:14
Status: success
Model: azure-o3
Total Steps: 6
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

# Expected language codes for English (US)
EXPECTED_LANGS = {"en-US", "en_US", "en-us"}


def _get_run_lang(run):
    """Return the language attribute value for a run if present, else None."""
    rPr = run._r.get_or_add_rPr()
    for attr_key, val in rPr.attrib.items():
        if attr_key.endswith('}lang') or attr_key == 'lang':
            return val
    return None


def _get_paragraph_lang(paragraph):
    """Return the language attribute value for a paragraph if present, else None."""
    pPr = paragraph._p.get_or_add_pPr()
    for attr_key, val in pPr.attrib.items():
        if attr_key.endswith('}lang') or attr_key == 'lang':
            return val
    return None


def verify_slide_title_language(file_path: str, slide_index: int = 55, expected_langs=EXPECTED_LANGS):
    """Verify that the title placeholder of a specific slide is set to English (US).

    Scoring breakdown:
        0.2 points – Title placeholder exists on the specified slide.
        0.3 points – At least one language attribute (run-level or paragraph-level) is set to English (US).
        0.5 points – ALL detected language attributes (run-level or paragraph-level) are set to English (US).

    Returns:
        float: total score between 0.0 and 1.0
    """
    total_score = 0.0

    # ---------- 1. Preliminary checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load PPTX file: {e}")
        return 0.0

    if len(prs.slides) <= slide_index:
        print(
            f"✗ Presentation has {len(prs.slides)} slides; slide {slide_index + 1} not available."
        )
        return 0.0

    slide = prs.slides[slide_index]

    # ---------- 2. Locate title placeholder (0.2 points) ----------
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            if shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            ):
                title_shape = shape
                break

    if title_shape is None:
        print("✗ No title placeholder found on slide 56 – cannot verify language.")
        return 0.0

    print("✓ Title placeholder found on slide 56 (0.2 points)")
    total_score += 0.2

    # ---------- 3. Collect language attributes ----------
    run_langs = []
    para_langs = []

    for paragraph in title_shape.text_frame.paragraphs:
        lang = _get_paragraph_lang(paragraph)
        if lang is not None:
            para_langs.append(lang)
        for run in paragraph.runs:
            lang = _get_run_lang(run)
            if lang is not None:
                run_langs.append(lang)

    combined_langs = run_langs if run_langs else para_langs  # Prefer run-level

    if not combined_langs:
        print(
            "✗ No language attributes found on title – spell-checker language not explicitly set."
        )
        print(f"Total score: {total_score}/1.0")
        return total_score

    # ---------- 4. Evaluate language attributes ----------
    print(f"Detected language attributes: {combined_langs}")

    any_correct = any(lang in expected_langs for lang in combined_langs)
    all_correct = all(lang in expected_langs for lang in combined_langs)

    if any_correct:
        print("✓ At least one language attribute correctly set to English (US) (0.3 points)")
        total_score += 0.3
        if all_correct:
            print("✓ All language attributes correctly set to English (US) (0.5 points)")
            total_score += 0.5
        else:
            incorrect = [lang for lang in combined_langs if lang not in expected_langs]
            print(f"✗ Some language attributes are incorrect: {incorrect}")
    else:
        print("✗ No language attribute set to English (US).")

    # ---------- 5. Final score ----------
    final_score = min(total_score, 1.0)
    print(f"Total score breakdown: {final_score}/1.0")
    return final_score


if __name__ == "__main__":
    FILE_PATH = (
        "/home/user/slide_56_is_driving_me_nutslibreoffice_impress_keeps_underlining_"
        "every_single_word_in_the_title_beca_golden.pptx"
    )
    reward = verify_slide_title_language(FILE_PATH)
    print(f"REWARD: {reward}")

