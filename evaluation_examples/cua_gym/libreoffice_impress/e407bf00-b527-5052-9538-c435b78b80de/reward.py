"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm working on a presentation and need to insert a new slide immediately after the title slide. I want it to have the 'Title and Content' format so I can add some key points. How do I do this in LibreOffice Impress?
Generated: 2025-08-07 10:31:19
Status: success
Model: o4-mini
Total Steps: 1
"""

import os
from pptx import Presentation

def verify_impress_task(file_path):
    print(f"Starting Impress verification for: {file_path}")
    total_score = 0.0
    
    # Requirement 1: File exists (0.2 points)
    try:
        if os.path.exists(file_path):
            print(f"✓ File exists: {file_path} (0.2 points)")
            total_score += 0.2
        else:
            print(f"✗ File not found: {file_path}")
            print(f"Final score: {total_score}")
            return min(total_score, 1.0)
    except Exception as e:
        print(f"✗ Error checking file existence: {e}")
        print(f"Final score: {total_score}")
        return min(total_score, 1.0)
    
    # Requirement 2: Load presentation and check slide count (0.3 points)
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully with {len(prs.slides)} slides")
        slide_count = len(prs.slides)
        print(f"Slide count found: {slide_count} (expected >= 2)")
        if slide_count >= 2:
            print(f"✓ Slide count requirement met (0.3 points)")
            total_score += 0.3
        else:
            print(f"✗ Slide count requirement not met (0 points)")
    except Exception as e:
        print(f"✗ Error loading presentation or checking slide count: {e}")
        print(f"Final score: {total_score}")
        return min(total_score, 1.0)
    
    # Debug: list all available slide layouts
    try:
        print("Available slide layouts in this presentation:")
        for i, layout in enumerate(prs.slide_layouts):
            name = layout.name if hasattr(layout, 'name') else '<no name>'
            print(f"  Layout {i}: '{name}'")
    except Exception as e:
        print(f"✗ Error listing slide layouts: {e}")
    
    # Requirement 3: Verify second slide uses 'Title and Content' layout (0.5 points)
    try:
        second_slide = prs.slides[1]
        layout_name = second_slide.slide_layout.name
        print(f"Second slide layout name: '{layout_name}'")
        if layout_name.lower().strip() == 'title and content':
            print(f"✓ Second slide layout is 'Title and Content' (0.5 points)")
            total_score += 0.5
        else:
            print(f"✗ Second slide layout is not 'Title and Content' (0 points)")
    except Exception as e:
        print(f"✗ Error verifying second slide layout: {e}")
    
    # Calculate final score
    final_score = min(total_score, 1.0)
    print(f"Final score: {final_score}")
    return final_score

# Execute verification and print reward
file_path = '/home/user/im_working_on_a_presentation_and_need_to_insert_a_new_slide_immediately_after_the_title_slide_i_want.pptx'
reward = verify_impress_task(file_path)
print(f"REWARD: {reward}")
