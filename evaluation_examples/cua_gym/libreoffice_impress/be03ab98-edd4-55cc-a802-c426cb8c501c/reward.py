"""
Reward Script: Apply text outline to title on slide 1
Task ID: impress_tct_078
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Text outline exists with correct width (1.5pt = 19050 EMU)
  Component 2 (0.30): Text outline color is dark blue (#1A237E)
  Component 3 (0.35): Text fill color is white (#FFFFFF)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_078'


def persist_app_state(domain: str):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slide 1 exists and has a title shape with text
    if len(prs.slides) < 1:
        print("FAIL: Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Find the title shape
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name and 'title' in shape.name.lower():
            title_shape = shape
            break

    # Fallback: try shapes.title
    if title_shape is None:
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass

    if title_shape is None or not title_shape.has_text_frame:
        print("FAIL: No title shape found on slide 1")
        print("REWARD: 0.0")
        return 0.0

    # Get runs from the title
    title_runs = []
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                title_runs.append(run)

    if not title_runs:
        print("FAIL: Title shape has no text runs")
        print("REWARD: 0.0")
        return 0.0

    # We check the first run (which contains 'Grand Opening')
    run = title_runs[0]
    rPr = run.font._element
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    # Component 1: Text outline exists with correct width 1.5pt = 19050 EMU (0.35 points)
    try:
        ln = rPr.find('a:ln', nsmap)
        if ln is not None:
            w_attr = ln.get('w')
            if w_attr is not None:
                w_val = int(w_attr)
                # 1.5pt = 19050 EMU; allow small tolerance (18000-20000)
                if 18000 <= w_val <= 20100:
                    print(f"PASS: Component 1 -- Text outline exists with width {w_val} EMU (~1.5pt) (0.35 pts)")
                    total_score += 0.35
                else:
                    print(f"FAIL: Component 1 -- Outline width is {w_val} EMU, expected ~19050 (1.5pt)")
            else:
                print("FAIL: Component 1 -- Outline element exists but no width attribute")
        else:
            print("FAIL: Component 1 -- No text outline (a:ln) element found on title run")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Text outline color is dark blue #1A237E (0.30 points)
    try:
        ln = rPr.find('a:ln', nsmap)
        if ln is not None:
            # Look for solidFill inside ln
            ln_fill = ln.find('a:solidFill', nsmap)
            if ln_fill is not None:
                srgb = ln_fill.find('a:srgbClr', nsmap)
                if srgb is not None:
                    color_val = srgb.get('val', '').upper()
                    if color_val == '1A237E':
                        print(f"PASS: Component 2 -- Outline color is #{color_val} (dark blue) (0.30 pts)")
                        total_score += 0.30
                    else:
                        print(f"FAIL: Component 2 -- Outline color is #{color_val}, expected #1A237E")
                else:
                    print("FAIL: Component 2 -- No srgbClr in outline fill")
            else:
                print("FAIL: Component 2 -- No solidFill inside outline element")
        else:
            print("FAIL: Component 2 -- No outline element, cannot check color")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Text fill color is white #FFFFFF (0.35 points)
    try:
        # Check via python-pptx API first
        color_rgb = None
        try:
            if run.font.color.type is not None:
                color_rgb = str(run.font.color.rgb).upper()
        except Exception:
            pass

        # Fallback: check XML solidFill directly on rPr
        if color_rgb is None:
            sf = rPr.find('a:solidFill', nsmap)
            if sf is not None:
                srgb = sf.find('a:srgbClr', nsmap)
                if srgb is not None:
                    color_rgb = srgb.get('val', '').upper()

        if color_rgb == 'FFFFFF':
            print(f"PASS: Component 3 -- Text fill color is #FFFFFF (white) (0.35 pts)")
            total_score += 0.35
        else:
            print(f"FAIL: Component 3 -- Text fill color is #{color_rgb}, expected #FFFFFF")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
