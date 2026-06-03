"""
Initial Setup: Create a 3-slide presentation with title 'Grand Opening' in solid black 44pt, no outline.
Task ID: impress_tct_078
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_078'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "Grand Opening"
    for para in title1.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # solid black

    # Subtitle
    subtitle = slide1.placeholders[1]
    subtitle.text = "Celebrating a New Chapter"
    for run in subtitle.text_frame.paragraphs[0].runs:
        run.font.name = "Arial"
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Event Agenda"
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.clear()
    agenda_items = [
        "9:00 AM - Welcome Remarks by CEO Maria Torres",
        "9:30 AM - Ribbon Cutting Ceremony",
        "10:00 AM - Guided Facility Tour",
        "11:00 AM - Product Showcase & Demonstrations",
        "12:00 PM - Networking Lunch & Refreshments",
        "1:30 PM - Panel Discussion: Future Innovations",
        "3:00 PM - Closing Remarks",
    ]
    for i, item in enumerate(agenda_items):
        if i == 0:
            tf2.paragraphs[0].text = item
        else:
            p = tf2.add_paragraph()
            p.text = item
            p.level = 0

    # --- Slide 3: Venue Details ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Venue & Contact"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.clear()
    details = [
        "Location: Harborview Convention Center, 425 Marina Blvd",
        "Date: Saturday, March 22, 2025",
        "Time: 9:00 AM - 4:00 PM",
        "RSVP: events@grandopeninghq.com",
        "Phone: (555) 842-1190",
        "Dress Code: Business Casual",
    ]
    for i, detail in enumerate(details):
        if i == 0:
            tf3.paragraphs[0].text = detail
        else:
            p = tf3.add_paragraph()
            p.text = detail
            p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
