"""
Reward Script: Add DRAFT watermark on every slide
Task ID: impress_stu_082
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): All 8 slides contain a shape with 'DRAFT' text
  Component 2 (0.25): Font size is 72pt on all watermark shapes
  Component 3 (0.20): Font color is #D0D0D0 (light gray) on all watermarks
  Component 4 (0.15): 50% opacity (alpha=50000) on all watermark text
  Component 5 (0.10): Rotation is approximately 45 degrees (315 deg in OOXML) on all watermarks
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_082'
EXPECTED_SLIDES = 8


def find_draft_shapes(slide):
    """Find all shapes containing 'DRAFT' text on a slide."""
    draft_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if full_text.upper() == 'DRAFT':
                draft_shapes.append(shape)
    return draft_shapes


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 8 slides
    num_slides = len(prs.slides)
    if num_slides != EXPECTED_SLIDES:
        print(f"PRECONDITION FAIL: Expected {EXPECTED_SLIDES} slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Collect DRAFT watermark shapes per slide
    watermarks = {}
    for idx, slide in enumerate(prs.slides):
        drafts = find_draft_shapes(slide)
        if drafts:
            watermarks[idx] = drafts[0]  # take first match

    # Component 1: All 8 slides have a DRAFT text shape (0.30 points)
    try:
        slides_with_draft = len(watermarks)
        if slides_with_draft == EXPECTED_SLIDES:
            print(f"PASS: Component 1 — All {EXPECTED_SLIDES} slides have DRAFT watermark (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 — Only {slides_with_draft}/{EXPECTED_SLIDES} slides have DRAFT watermark")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # If no watermarks found at all, remaining checks are moot
    if not watermarks:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Font size is 72pt on all watermarks (0.25 points)
    try:
        correct_size_count = 0
        for idx, shape in watermarks.items():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        # 72pt = 914400 EMU
                        size_pt = run.font.size / 12700
                        if abs(size_pt - 72.0) < 1.0:  # tolerance of 1pt
                            correct_size_count += 1
                            break
                else:
                    continue
                break

        if correct_size_count == len(watermarks):
            print(f"PASS: Component 2 — All watermarks have 72pt font size (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — {correct_size_count}/{len(watermarks)} watermarks have 72pt font")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Font color is #D0D0D0 on all watermarks (0.20 points)
    try:
        correct_color_count = 0
        for idx, shape in watermarks.items():
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == 'D0D0D0':
                                correct_color_count += 1
                                break
                    except Exception:
                        pass
                else:
                    continue
                break

        if correct_color_count == len(watermarks):
            print(f"PASS: Component 3 — All watermarks have #D0D0D0 color (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — {correct_color_count}/{len(watermarks)} watermarks have #D0D0D0 color")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 50% opacity / alpha=50000 on all watermarks (0.15 points)
    try:
        correct_alpha_count = 0
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for idx, shape in watermarks.items():
            # Parse shape XML to find alpha value
            sp_el = shape._element
            alpha_els = sp_el.findall(f'.//{{{ns_a}}}alpha')
            alpha_vals = [int(el.get('val', '0')) for el in alpha_els if el.get('val')]
            # 50% opacity = alpha 50000 (in 1/1000ths of percent)
            if any(45000 <= v <= 55000 for v in alpha_vals):
                correct_alpha_count += 1
            else:
                print(f"  Slide {idx+1}: no valid alpha found (values: {alpha_vals})")

        if correct_alpha_count == len(watermarks):
            print(f"PASS: Component 4 — All watermarks have ~50% opacity (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — {correct_alpha_count}/{len(watermarks)} watermarks have 50% opacity")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Rotation is ~45 degrees (315 in python-pptx = 18900000 in OOXML) (0.10 points)
    try:
        correct_rotation_count = 0
        for idx, shape in watermarks.items():
            rot = shape.rotation
            # 315 degrees (python-pptx) = -45 degrees = 45 degree diagonal
            # Also accept 45 directly or -45
            # Normalize to 0-360 range
            norm_rot = rot % 360
            if abs(norm_rot - 315.0) < 5.0 or abs(norm_rot - 45.0) < 5.0:
                correct_rotation_count += 1
            else:
                print(f"  Slide {idx+1}: rotation is {rot} (normalized {norm_rot})")

        if correct_rotation_count == len(watermarks):
            print(f"PASS: Component 5 — All watermarks have ~45-degree rotation (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — {correct_rotation_count}/{len(watermarks)} watermarks have correct rotation")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
