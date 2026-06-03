"""
Initial Setup: Essay Draft presentation with 8 slides (no watermark)
Task ID: impress_stu_082
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_082'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "The Impact of Renewable Energy on Global Economics"
    slide1.placeholders[1].text = "A Comprehensive Analysis\nPrepared by: Elena Rodriguez\nEnvironmental Studies 401 — Spring 2026"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "The global energy landscape is undergoing a profound transformation."
    p2 = body2.add_paragraph()
    p2.text = "Renewable energy sources now account for 29% of global electricity generation as of 2025, up from just 20% a decade ago."
    p3 = body2.add_paragraph()
    p3.text = "This presentation examines the economic implications of this shift across developed and developing nations."
    p4 = body2.add_paragraph()
    p4.text = "Key areas of focus include job creation, investment trends, and the decline of fossil fuel industries."

    # --- Slide 3: Solar Energy Growth ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Solar Energy: Leading the Charge"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Solar capacity has grown at a compound annual rate of 26% since 2015."
    for text in [
        "Global installed capacity reached 1,580 GW by end of 2025",
        "Manufacturing costs dropped 89% over the past 15 years",
        "China, India, and the United States remain the top three markets",
        "Utility-scale solar is now cheaper than coal in most regions",
    ]:
        p = body3.add_paragraph()
        p.text = text

    # --- Slide 4: Wind Energy Developments ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Wind Energy: Offshore Expansion"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Offshore wind has emerged as a critical growth sector."
    for text in [
        "Total global offshore wind capacity surpassed 75 GW in 2025",
        "The North Sea and East Asia dominate current installations",
        "Floating wind platforms are opening previously inaccessible deep-water sites",
        "Average turbine size has increased from 3 MW to 15 MW in two decades",
        "Levelized cost of offshore wind fell below $70/MWh",
    ]:
        p = body4.add_paragraph()
        p.text = text

    # --- Slide 5: Economic Impact ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Economic Impact on Employment"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "The renewable energy sector employed over 13.7 million people globally in 2025."
    for text in [
        "Solar photovoltaics: 4.9 million jobs worldwide",
        "Wind energy: 1.4 million jobs, with offshore growing fastest",
        "Bioenergy and hydropower: 5.2 million jobs combined",
        "Net job creation exceeds losses in fossil fuel sectors by 3:1 ratio",
        "Emerging economies in Southeast Asia see the highest growth rates",
    ]:
        p = body5.add_paragraph()
        p.text = text

    # --- Slide 6: Investment Trends ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Global Investment Trends"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Annual renewable energy investment exceeded $500 billion in 2025."
    for text in [
        "Private sector capital accounts for 72% of total investment",
        "Green bond issuance reached $620 billion, a new record",
        "Venture capital in energy storage startups tripled since 2022",
        "Government subsidies are gradually shifting from production to R&D",
    ]:
        p = body6.add_paragraph()
        p.text = text

    # --- Slide 7: Challenges and Barriers ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Challenges and Barriers"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Despite rapid growth, significant challenges remain."
    for text in [
        "Grid infrastructure requires $4.5 trillion in upgrades by 2040",
        "Intermittency and storage limitations persist for solar and wind",
        "Supply chain concentration: 80% of solar panels manufactured in China",
        "Permitting delays average 4.5 years for large onshore wind projects in Europe",
        "Critical mineral supply (lithium, cobalt, rare earths) faces geopolitical risks",
    ]:
        p = body7.add_paragraph()
        p.text = text

    # --- Slide 8: Conclusion ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Conclusion and Outlook"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "The transition to renewable energy is irreversible and accelerating."
    for text in [
        "Policy frameworks like the EU Green Deal and US IRA are driving adoption",
        "Cost parity with fossil fuels has been achieved in most markets",
        "Energy independence is becoming a national security priority",
        "The economic benefits increasingly outweigh transition costs",
        "By 2035, renewables are projected to generate over 60% of global electricity",
    ]:
        p = body8.add_paragraph()
        p.text = text

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
