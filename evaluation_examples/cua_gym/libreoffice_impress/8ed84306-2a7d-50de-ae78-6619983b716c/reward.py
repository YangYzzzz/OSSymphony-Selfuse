"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 9 I’ve got “Picture 1” sitting there plain as day. In LibreOffice Impress, how can I ring that image with a #000000 dashed outline that’s exactly 1.00 pt wide?
Generated: 2025-09-10 21:45:05
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

# Path to the presentation created/edited by the task performer
FILE_PATH = '/home/user/on_slide_9_ive_got_picture_1_sitting_there_plain_as_day_in_libreoffice_impress_how_can_i_ring_that_i_golden.pptx'


def verify_dashed_outline(file_path: str) -> float:
    """Verify that on slide 9 a picture has a 1 pt (#000000) dashed outline.

    Scoring (progressive, totals 1.0):
        0.25 – A picture exists on slide 9
        0.25 – Outline colour is exactly #000000
        0.25 – Outline dash style is DASH
        0.25 – Outline width equals 1 pt (≈ 12700 EMU)
    """

    print(f"Verifying presentation: {file_path}")
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0

    # Load presentation (no points for merely loading)
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load PPTX: {e}")
        return 0.0

    if len(prs.slides) < 9:
        print("✗ Presentation has fewer than 9 slides")
        return 0.0

    slide = prs.slides[8]  # zero-based index => slide 9

    score = 0.0
    weights = {
        'picture_found': 0.25,
        'color_correct': 0.25,
        'dash_correct': 0.25,
        'width_correct': 0.25,
    }

    # Look for first picture on slide 9
    picture = next((sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE), None)

    if not picture:
        print("✗ No picture found on slide 9")
        return 0.0  # nothing else to check

    print(f"✓ Picture found: {picture.name}")
    score += weights['picture_found']

    # Evaluate the outline/line properties
    ln = picture.line  # pptx always gives a LineFormat object

    # 1) Colour check – must be pure black (#000000)
    colour_ok = False
    if ln.color is not None and ln.color.rgb is not None:
        rgb = str(ln.color.rgb).upper()
        print(f"  Line colour RGB: {rgb}")
        if rgb == '000000':
            colour_ok = True
            score += weights['color_correct']
            print("  ✓ Colour is #000000")
        else:
            print("  ✗ Colour is not #000000")
    else:
        print("  ✗ No RGB colour defined for outline")

    # 2) Dash style check – must be DASH
    dash_ok = False
    if ln.dash_style is not None:
        print(f"  Dash style: {ln.dash_style}")
        if ln.dash_style == MSO_LINE_DASH_STYLE.DASH:
            dash_ok = True
            score += weights['dash_correct']
            print("  ✓ Dash style is DASH")
        else:
            print("  ✗ Dash style is not DASH")
    else:
        print("  ✗ No dash style set")

    # 3) Width check – must be 1 pt (≈ 12700 EMU). Allow ±100 EMU tolerance.
    width_ok = False
    if ln.width is not None:
        width_emu = ln.width
        print(f"  Line width (EMU): {width_emu}")
        if abs(width_emu - 12700) <= 100:
            width_ok = True
            score += weights['width_correct']
            print("  ✓ Width is approximately 1 pt")
        else:
            print("  ✗ Width is not 1 pt")
    else:
        print("  ✗ Line width unavailable")

    # Final reporting
    print(f"Total score breakdown: {score}/1.0")
    final_score = min(score, 1.0)
    return final_score


# Execute verification and print reward
reward = verify_dashed_outline(FILE_PATH)
print(f"REWARD: {reward}")

