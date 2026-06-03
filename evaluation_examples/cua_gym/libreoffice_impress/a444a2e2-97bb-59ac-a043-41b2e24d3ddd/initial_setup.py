"""
Initial Setup: Create a 10-slide keynote presentation with slide 2 titled 'Inspiration'
Task ID: impress_rp_028
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_028'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Keynote Address"
    slide1.placeholders[1].text = "Annual Leadership Summit 2025"

    # --- Slide 2: Inspiration (title only, empty body) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add title as a text box at top
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Inspiration"
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    # No other shapes - this is the empty body area for the agent to fill

    # --- Slide 3: Vision & Mission ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Vision & Mission"
    p3.runs[0].font.size = Pt(36)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b3 = body3.text_frame
    tf_b3.word_wrap = True
    tf_b3.paragraphs[0].text = "Our vision is to empower organizations worldwide through transformative technology solutions that drive sustainable growth and meaningful impact."
    p_m = tf_b3.add_paragraph()
    p_m.text = "Mission: Deliver innovative platforms that simplify complex workflows, foster collaboration, and accelerate digital transformation across every industry."
    for para in tf_b3.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 4: Market Landscape ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Market Landscape"
    p4.runs[0].font.size = Pt(36)
    p4.runs[0].font.bold = True
    p4.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b4 = body4.text_frame
    tf_b4.word_wrap = True
    bullets4 = [
        "Global SaaS market projected to reach $908B by 2030",
        "Enterprise AI adoption increased 47% year-over-year",
        "Remote collaboration tools grew 312% since 2020",
        "Cybersecurity spending expected to surpass $300B by 2027",
        "Cloud-native architectures now power 78% of new deployments",
    ]
    tf_b4.paragraphs[0].text = bullets4[0]
    for b in bullets4[1:]:
        p_b = tf_b4.add_paragraph()
        p_b.text = b
    for para in tf_b4.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 5: Key Achievements ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    tb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf5 = tb5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Achievements"
    p5.runs[0].font.size = Pt(36)
    p5.runs[0].font.bold = True
    p5.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b5 = body5.text_frame
    tf_b5.word_wrap = True
    achievements = [
        "Revenue growth of 34% reaching $2.8 billion in FY2024",
        "Expanded to 42 countries with 15,000+ enterprise clients",
        "Launched 3 breakthrough products: DataSync Pro, CloudBridge, InsightAI",
        "Customer satisfaction score improved to 94.2% (up from 88.7%)",
        "Reduced carbon footprint by 28% through green data center initiatives",
    ]
    tf_b5.paragraphs[0].text = achievements[0]
    for a in achievements[1:]:
        p_a = tf_b5.add_paragraph()
        p_a.text = a
    for para in tf_b5.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 6: Technology Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    tb6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf6 = tb6.text_frame
    p6 = tf6.paragraphs[0]
    p6.text = "Technology Roadmap"
    p6.runs[0].font.size = Pt(36)
    p6.runs[0].font.bold = True
    p6.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b6 = body6.text_frame
    tf_b6.word_wrap = True
    roadmap = [
        "Q1 2025: Launch AI-powered analytics dashboard for enterprise clients",
        "Q2 2025: Roll out edge computing framework for IoT integration",
        "Q3 2025: Release cross-platform mobile SDK with offline-first architecture",
        "Q4 2025: Deploy federated learning infrastructure for privacy-preserving AI",
    ]
    tf_b6.paragraphs[0].text = roadmap[0]
    for item in roadmap[1:]:
        p_r = tf_b6.add_paragraph()
        p_r.text = item
    for para in tf_b6.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 7: Team & Culture ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    tb7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf7 = tb7.text_frame
    p7 = tf7.paragraphs[0]
    p7.text = "Team & Culture"
    p7.runs[0].font.size = Pt(36)
    p7.runs[0].font.bold = True
    p7.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b7 = body7.text_frame
    tf_b7.word_wrap = True
    tf_b7.paragraphs[0].text = "Our team of 8,200 professionals across 42 countries represents the heartbeat of our innovation engine. We cultivate a culture of curiosity, inclusion, and bold experimentation."
    p_t = tf_b7.add_paragraph()
    p_t.text = "Employee engagement score: 91% | Voluntary turnover: 6.2% | Internal mobility: 34%"
    for para in tf_b7.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 8: Strategic Partnerships ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    tb8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf8 = tb8.text_frame
    p8 = tf8.paragraphs[0]
    p8.text = "Strategic Partnerships"
    p8.runs[0].font.size = Pt(36)
    p8.runs[0].font.bold = True
    p8.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b8 = body8.text_frame
    tf_b8.word_wrap = True
    partners = [
        "Microsoft Azure: Co-developed hybrid cloud solutions for Fortune 500 clients",
        "AWS: Joint AI/ML certification program reaching 50,000 developers",
        "Siemens: Industrial IoT integration for smart manufacturing",
        "Accenture: Global digital transformation consulting alliance",
    ]
    tf_b8.paragraphs[0].text = partners[0]
    for pt in partners[1:]:
        p_pt = tf_b8.add_paragraph()
        p_pt.text = pt
    for para in tf_b8.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 9: Financial Outlook ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    tb9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf9 = tb9.text_frame
    p9 = tf9.paragraphs[0]
    p9.text = "Financial Outlook"
    p9.runs[0].font.size = Pt(36)
    p9.runs[0].font.bold = True
    p9.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    body9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_b9 = body9.text_frame
    tf_b9.word_wrap = True
    outlook = [
        "Projected revenue for FY2025: $3.6 billion (29% growth)",
        "Operating margin target: 22-24%",
        "R&D investment increasing to 18% of revenue",
        "Free cash flow expected to exceed $800 million",
        "Dividend increase of 15% planned for Q3 2025",
    ]
    tf_b9.paragraphs[0].text = outlook[0]
    for o in outlook[1:]:
        p_o = tf_b9.add_paragraph()
        p_o.text = o
    for para in tf_b9.paragraphs:
        for r in para.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    tb10 = slide10.shapes.add_textbox(Inches(3), Inches(2.5), Inches(7), Inches(2))
    tf10 = tb10.text_frame
    p10 = tf10.paragraphs[0]
    p10.text = "Thank You"
    p10.alignment = PP_ALIGN.CENTER
    p10.runs[0].font.size = Pt(48)
    p10.runs[0].font.bold = True
    p10.runs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p10_sub = tf10.add_paragraph()
    p10_sub.text = "Questions & Discussion"
    p10_sub.alignment = PP_ALIGN.CENTER
    p10_sub.runs[0].font.size = Pt(24)
    p10_sub.runs[0].font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
