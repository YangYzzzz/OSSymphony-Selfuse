"""
Reward Script: Pull-quote text box with border bar on slide 2
Task ID: impress_rp_028
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Opening quote mark paragraph - 72pt light gray
  Component 2 (0.3): Quote text paragraph - 24pt italic dark gray
  Component 3 (0.2): Attribution paragraph - 14pt right-aligned
  Component 4 (0.3): Left border bar rectangle - 4pt wide, E74C3C fill
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_028'


def find_pullquote_textbox(slide):
    """Find a text box on slide 2 that contains the pull-quote content.
    Must have at least 3 paragraphs and NOT be the existing title/Inspiration textbox."""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        # Collect all non-empty paragraph texts
        nonempty = [p for p in tf.paragraphs if p.text.strip()]
        if len(nonempty) < 2:
            continue
        full_text = " ".join(p.text.strip() for p in tf.paragraphs)
        # Must contain part of the quote or the attribution
        if "Innovation" in full_text or "Steve Jobs" in full_text:
            candidates.append(shape)
    return candidates[0] if candidates else None


def find_border_bar(slide, textbox_shape):
    """Find a rectangle shape that acts as a left border bar near the textbox."""
    if textbox_shape is None:
        return None
    tb_left = textbox_shape.left
    tb_top = textbox_shape.top
    tb_height = textbox_shape.height

    for shape in slide.shapes:
        # Look for auto shapes (rectangles)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
            # Must be narrow (border bar) - width should be small relative to height
            if shape.width < Emu(200000) and shape.height > Emu(500000):
                return shape
    return None


def color_matches(rgb_val, expected_hex, tolerance=30):
    """Check if an RGB color matches expected hex within tolerance."""
    if rgb_val is None:
        return False
    actual = str(rgb_val).upper()
    expected = expected_hex.upper()
    if actual == expected:
        return True
    # Parse and compare with tolerance
    try:
        ar, ag, ab = int(actual[0:2], 16), int(actual[2:4], 16), int(actual[4:6], 16)
        er, eg, eb = int(expected[0:2], 16), int(expected[2:4], 16), int(expected[4:6], 16)
        return abs(ar - er) <= tolerance and abs(ag - eg) <= tolerance and abs(ab - eb) <= tolerance
    except (ValueError, IndexError):
        return False


def size_matches(actual_size, expected_pt, tolerance_pt=4):
    """Check if font size (in EMU) is approximately expected pt value."""
    if actual_size is None:
        return False
    actual_pt = actual_size / 12700
    return abs(actual_pt - expected_pt) <= tolerance_pt


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 2 slides
    if len(prs.slides) < 2:
        print("FAIL: Presentation has fewer than 2 slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # Find the pull-quote text box
    quote_tb = find_pullquote_textbox(slide)

    # Component 1: Opening quote mark paragraph (0.2 points)
    # Must have a paragraph with a large quotation mark (~72pt) in light gray
    try:
        if quote_tb is None:
            print("FAIL: Component 1 -- No pull-quote text box found on slide 2")
        else:
            tf = quote_tb.text_frame
            paras = tf.paragraphs
            found_quote_mark = False
            for para in paras:
                runs = [r for r in para.runs if r.text.strip()]
                if not runs:
                    continue
                first_run = runs[0]
                text = first_run.text.strip()
                # Check for quotation mark character(s)
                if any(c in text for c in ['\u201c', '\u201d', '"', '\u0022', '\u00ab', '\u2018', '\u2019']):
                    # Check size is approximately 72pt (914400 EMU)
                    if size_matches(first_run.font.size, 72, tolerance_pt=8):
                        # Check color is light gray (~C0C0C0)
                        try:
                            rgb = first_run.font.color.rgb if first_run.font.color.type is not None else None
                            if color_matches(rgb, "C0C0C0", tolerance=50):
                                found_quote_mark = True
                                print(f"PASS: Component 1 -- Quote mark '{text}' in {first_run.font.size/12700}pt, color={rgb} (0.2 pts)")
                                total_score += 0.2
                                break
                            else:
                                print(f"FAIL: Component 1 -- Quote mark color mismatch: expected ~C0C0C0, found {rgb}")
                        except Exception:
                            print(f"FAIL: Component 1 -- Could not read quote mark color")
                    else:
                        actual_pt = first_run.font.size / 12700 if first_run.font.size else "None"
                        print(f"FAIL: Component 1 -- Quote mark size mismatch: expected ~72pt, found {actual_pt}pt")
            if not found_quote_mark and quote_tb is not None:
                # Check if we just didn't print a fail message yet
                all_texts = [p.text.strip() for p in paras if p.text.strip()]
                print(f"FAIL: Component 1 -- No quote mark paragraph found. Paragraphs: {all_texts[:5]}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Quote text in 24pt italic dark gray (0.3 points)
    # Must have a paragraph with the quote text in ~24pt, italic, color ~333333
    try:
        if quote_tb is None:
            print("FAIL: Component 2 -- No pull-quote text box found")
        else:
            tf = quote_tb.text_frame
            found_quote_text = False
            for para in tf.paragraphs:
                runs = [r for r in para.runs if r.text.strip()]
                if not runs:
                    continue
                full_para_text = "".join(r.text for r in runs).strip()
                # Check for the quote text
                if "Innovation" in full_para_text and "leader" in full_para_text and "follower" in full_para_text:
                    sub_score = 0.0
                    first_run = runs[0]

                    # Check size ~24pt
                    if size_matches(first_run.font.size, 24, tolerance_pt=4):
                        sub_score += 0.1
                        print(f"  PASS: Quote text size ~24pt (actual: {first_run.font.size/12700}pt)")
                    else:
                        actual_pt = first_run.font.size / 12700 if first_run.font.size else "None"
                        print(f"  FAIL: Quote text size expected ~24pt, found {actual_pt}pt")

                    # Check italic
                    is_italic = first_run.font.italic
                    if is_italic is True:
                        sub_score += 0.1
                        print(f"  PASS: Quote text is italic")
                    else:
                        print(f"  FAIL: Quote text italic={is_italic}, expected True")

                    # Check color ~333333
                    try:
                        rgb = first_run.font.color.rgb if first_run.font.color.type is not None else None
                        if color_matches(rgb, "333333", tolerance=30):
                            sub_score += 0.1
                            print(f"  PASS: Quote text color={rgb} (~333333)")
                        else:
                            print(f"  FAIL: Quote text color expected ~333333, found {rgb}")
                    except Exception:
                        print(f"  FAIL: Could not read quote text color")

                    total_score += sub_score
                    found_quote_text = True
                    pts = sub_score
                    print(f"RESULT: Component 2 -- {pts}/0.3 pts")
                    break

            if not found_quote_text:
                print("FAIL: Component 2 -- Quote text not found in any paragraph")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Attribution '- Steve Jobs' in 14pt right-aligned (0.2 points)
    try:
        if quote_tb is None:
            print("FAIL: Component 3 -- No pull-quote text box found")
        else:
            tf = quote_tb.text_frame
            found_attribution = False
            for para in tf.paragraphs:
                runs = [r for r in para.runs if r.text.strip()]
                if not runs:
                    continue
                full_para_text = "".join(r.text for r in runs).strip()
                if "Steve Jobs" in full_para_text:
                    sub_score = 0.0
                    first_run = runs[0]

                    # Check size ~14pt
                    if size_matches(first_run.font.size, 14, tolerance_pt=4):
                        sub_score += 0.1
                        print(f"  PASS: Attribution size ~14pt (actual: {first_run.font.size/12700}pt)")
                    else:
                        actual_pt = first_run.font.size / 12700 if first_run.font.size else "None"
                        print(f"  FAIL: Attribution size expected ~14pt, found {actual_pt}pt")

                    # Check right alignment
                    alignment = para.alignment
                    if alignment == PP_ALIGN.RIGHT:
                        sub_score += 0.1
                        print(f"  PASS: Attribution is right-aligned")
                    else:
                        print(f"  FAIL: Attribution alignment={alignment}, expected RIGHT")

                    total_score += sub_score
                    found_attribution = True
                    print(f"RESULT: Component 3 -- {sub_score}/0.2 pts")
                    break

            if not found_attribution:
                print("FAIL: Component 3 -- Attribution '- Steve Jobs' not found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Left border bar rectangle (0.3 points)
    # Must be a narrow rectangle with E74C3C fill, positioned at the left edge of the quote text box
    try:
        border_bar = find_border_bar(slide, quote_tb)
        if border_bar is None:
            print("FAIL: Component 4 -- No border bar rectangle found on slide 2")
        else:
            sub_score = 0.0

            # Check fill color is E74C3C
            try:
                fill = border_bar.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    fill_rgb = fill.fore_color.rgb
                    if color_matches(fill_rgb, "E74C3C", tolerance=20):
                        sub_score += 0.15
                        print(f"  PASS: Border bar fill color={fill_rgb} (~E74C3C)")
                    else:
                        print(f"  FAIL: Border bar fill color={fill_rgb}, expected ~E74C3C")
                else:
                    print(f"  FAIL: Border bar fill type={fill.type}, expected SOLID (1)")
            except Exception as e:
                print(f"  FAIL: Border bar fill check error: {e}")

            # Check width is approximately 4pt (50800 EMU)
            # 4pt width for a shape means ~50800 EMU (4 * 12700)
            bar_width = border_bar.width
            if bar_width <= Emu(200000):  # Narrow enough to be a border bar
                sub_score += 0.075
                print(f"  PASS: Border bar width={bar_width} EMU (~{bar_width/12700:.1f}pt) -- narrow bar")
            else:
                print(f"  FAIL: Border bar width={bar_width} EMU -- too wide for border bar")

            # Check height approximately matches textbox height
            if quote_tb is not None:
                tb_height = quote_tb.height
                bar_height = border_bar.height
                if tb_height > 0:
                    ratio = bar_height / tb_height
                    if 0.5 <= ratio <= 1.5:
                        sub_score += 0.075
                        print(f"  PASS: Border bar height matches textbox (ratio={ratio:.2f})")
                    else:
                        print(f"  FAIL: Border bar height ratio={ratio:.2f}, expected ~1.0")

            total_score += sub_score
            print(f"RESULT: Component 4 -- {sub_score}/0.3 pts")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(round(total_score, 4), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
