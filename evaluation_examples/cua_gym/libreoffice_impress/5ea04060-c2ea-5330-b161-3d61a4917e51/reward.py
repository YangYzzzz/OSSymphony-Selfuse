"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 145, the heading is still in black and looks disconnected from the image. Grab the main blue hue from Picture 1 with the eyedropper (it reads as #1E90FF in RGB) and apply that exact color to the title text.
Generated: 2025-09-10 17:42:37
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER

"""
Reward Script for Task:
"On slide 145, change the heading colour to the exact blue hue #1E90FF (RGB 30,144,255)."

Verification Logic:
1. Load the provided presentation file (no points for just loading!).
2. Ensure slide 145 exists.
3. Locate the heading (priority: title placeholder; fallback: upper-most text box with non-empty text).
4. Inspect every run in the heading and compare its RGB value to the required colour.
5. Scoring (progressive):
   • 1.0 – every run in the heading is the exact blue.
   • 0.5 – at least one run in the heading is blue but not all.
   • 0.2 – heading untouched, but some other text on that slide is blue (shows partial effort).
   • 0.0 – no evidence of the colour on the slide / heading not found.

Prints detailed diagnostics and finally outputs:  "REWARD: X.X"
"""

def get_heading_shape(slide):
    """Return the heading shape on a slide.
    1) Try title placeholder.
    2) Otherwise, pick the upper-most text shape that contains text."""
    for shape in slide.shapes:
        try:
            if (shape.is_placeholder and 
                shape.placeholder_format.type == PP_PLACEHOLDER.TITLE and 
                shape.has_text_frame and shape.text.strip()):
                return shape
        except Exception:
            # Some shapes may not expose placeholder_format; ignore safely
            pass

    heading_shape = None
    min_top = None
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        if not shape.text.strip():
            continue
        if min_top is None or shape.top < min_top:
            min_top = shape.top
            heading_shape = shape
    return heading_shape


def verify_heading_color(file_path):
    TARGET_RGB = RGBColor(0x1E, 0x90, 0xFF)  # #1E90FF

    # --- prerequisite checks: file must exist & load ---
    if not os.path.exists(file_path):
        print('✗ File does not exist:', file_path)
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print('✗ Could not load PPTX:', e)
        return 0.0  # cannot verify anything further

    # --- ensure slide 145 exists ---
    slide_index = 144  # 0-based index
    if slide_index >= len(prs.slides):
        print(f'✗ Slide 145 not found (presentation has {len(prs.slides)} slides)')
        return 0.0

    slide = prs.slides[slide_index]

    # --- locate heading shape ---
    heading_shape = get_heading_shape(slide)
    if heading_shape is None:
        print('✗ Could not locate heading text on slide 145')
        return 0.0
    print('✓ Found heading text:', heading_shape.text.strip())

    # --- analyse heading run colours ---
    total_runs = 0
    blue_runs = 0
    for paragraph in heading_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            total_runs += 1
            rgb = run.font.color.rgb  # returns None if colour is theme/auto
            if rgb == TARGET_RGB:
                blue_runs += 1
    print(f'Heading analysis: {blue_runs}/{total_runs} runs are target blue')

    # --- calculate score ---
    score = 0.0
    if total_runs == 0:
        print('✗ Heading contains no runs – unable to verify colour')
    elif blue_runs == total_runs:
        print('✓ All heading text correctly coloured #1E90FF')
        score = 1.0
    elif blue_runs > 0:
        print('⚠ Partial success – some, but not all, heading text coloured')
        score = 0.5
    else:
        # check if any other text on the slide has the correct colour – tiny partial credit
        any_blue_elsewhere = False
        for shape in slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.color.rgb == TARGET_RGB:
                        any_blue_elsewhere = True
                        break
                if any_blue_elsewhere:
                    break
            if any_blue_elsewhere:
                break
        if any_blue_elsewhere:
            print('⚠ Found the correct blue colour on the slide, but not on heading')
            score = 0.2
        else:
            print('✗ Required blue colour not found on slide 145')
            score = 0.0

    return score

# -----------------------------------------------------------------------------
# MAIN EXECUTION
# -----------------------------------------------------------------------------
FILE_PATH = "/home/user/on_slide_145_the_heading_is_still_in_black_and_looks_disconnected_from_the_image_grab_the_main_blue__golden.pptx"
reward = verify_heading_color(FILE_PATH)
print(f"REWARD: {reward}")
