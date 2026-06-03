"""
Reward Script: Infographic-style slide with four colored circles containing statistics
Task ID: impress_stu_056
Domain: libreoffice_impress
Scoring:
  C1: Title text "Key Findings at a Glance" (0.10)
  C2: Four oval/circle auto-shapes on slide 3 (0.15)
  C3: Correct statistic text in each circle (0.25)
  C4: Correct fill colors on circles (0.20)
  C5: Correct labels below circles (0.15)
  C6: Text formatting in circles — 36pt bold white (0.15)
"""

import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_056'

# Expected data for the four circles
EXPECTED_STATS = {
    '87%': '27AE60',       # green
    '1,200+': '2980B9',    # blue
    '4.2/5.0': 'F39C12',   # orange
    '92%': 'C0392B',       # red
}

EXPECTED_LABELS = {
    '87%': 'Satisfaction Rate',
    '1,200+': 'Respondents',
    '4.2/5.0': 'Average Rating',
    '92%': 'Would Recommend',
}


def get_shape_text(shape):
    """Extract concatenated text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    return "".join(para.text for para in shape.text_frame.paragraphs).strip()


def get_auto_shapes(slide):
    """Get all AUTO_SHAPE type shapes from a slide."""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


def get_text_boxes(slide):
    """Get all TEXT_BOX type shapes from a slide."""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]  # 0-indexed

    # Component 1: Title text "Key Findings at a Glance" (0.10 points)
    try:
        all_texts = [get_shape_text(s) for s in slide3.shapes if s.has_text_frame]
        title_matches = [t for t in all_texts if 'key findings at a glance' in t.lower()]
        if len(title_matches) > 0:
            print(f"PASS: Component 1 — Title 'Key Findings at a Glance' found (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — 'Key Findings at a Glance' not found. Texts: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Four oval/circle auto-shapes on slide 3 (0.15 points)
    try:
        auto_shapes = get_auto_shapes(slide3)
        num_ovals = len(auto_shapes)
        if num_ovals >= 4:
            print(f"PASS: Component 2 — Found {num_ovals} auto-shapes (need >= 4) (0.15 pts)")
            total_score += 0.15
        elif num_ovals >= 2:
            partial = 0.15 * (num_ovals / 4.0)
            print(f"PARTIAL: Component 2 — Found {num_ovals}/4 auto-shapes ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Found {num_ovals} auto-shapes, need 4")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct statistic text in each circle (0.25 points)
    # Each correct stat = 0.0625 points
    try:
        auto_shapes = get_auto_shapes(slide3)
        found_stats = set()
        for shape in auto_shapes:
            text = get_shape_text(shape).strip()
            if text in EXPECTED_STATS:
                found_stats.add(text)

        stat_score = len(found_stats) * (0.25 / 4.0)
        if len(found_stats) == 4:
            print(f"PASS: Component 3 — All 4 statistics found: {found_stats} (0.25 pts)")
            total_score += 0.25
        elif len(found_stats) > 0:
            print(f"PARTIAL: Component 3 — Found {len(found_stats)}/4 statistics: {found_stats} ({stat_score:.4f} pts)")
            total_score += stat_score
        else:
            texts = [get_shape_text(s) for s in auto_shapes]
            print(f"FAIL: Component 3 — No expected statistics found. Auto-shape texts: {texts}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Correct fill colors on circles (0.20 points)
    # Each correct color = 0.05 points
    try:
        auto_shapes = get_auto_shapes(slide3)
        color_matches = 0
        for shape in auto_shapes:
            text = get_shape_text(shape).strip()
            if text in EXPECTED_STATS:
                expected_color = EXPECTED_STATS[text].upper()
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID
                        actual_color = str(fill.fore_color.rgb).upper()
                        if actual_color == expected_color:
                            color_matches += 1
                            print(f"  Color OK: '{text}' has fill {actual_color}")
                        else:
                            print(f"  Color MISMATCH: '{text}' expected {expected_color}, got {actual_color}")
                    else:
                        print(f"  Color FAIL: '{text}' fill type is {fill.type}, not solid")
                except Exception as e:
                    print(f"  Color ERROR for '{text}': {e}")

        color_score = color_matches * (0.20 / 4.0)
        if color_matches == 4:
            print(f"PASS: Component 4 — All 4 circle colors correct (0.20 pts)")
            total_score += 0.20
        elif color_matches > 0:
            print(f"PARTIAL: Component 4 — {color_matches}/4 colors correct ({color_score:.4f} pts)")
            total_score += color_score
        else:
            print(f"FAIL: Component 4 — No correct fill colors found")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Correct labels below circles (0.15 points)
    # Each correct label = 0.0375 points
    try:
        text_boxes = get_text_boxes(slide3)
        # Collect label texts from text boxes (excluding the title textbox)
        label_texts = []
        for tb in text_boxes:
            t = get_shape_text(tb).strip()
            if t and 'key findings' not in t.lower():
                label_texts.append(t)

        expected_labels_set = set(EXPECTED_LABELS.values())
        found_labels = set()
        for lt in label_texts:
            if lt in expected_labels_set:
                found_labels.add(lt)

        label_count = len(found_labels)
        label_score = label_count * (0.15 / 4.0)
        if label_count == 4:
            print(f"PASS: Component 5 — All 4 labels found: {found_labels} (0.15 pts)")
            total_score += 0.15
        elif label_count > 0:
            print(f"PARTIAL: Component 5 — {label_count}/4 labels found: {found_labels} ({label_score:.4f} pts)")
            total_score += label_score
        else:
            print(f"FAIL: Component 5 — No expected labels found. Text boxes: {label_texts}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Text formatting in circles — 36pt bold white (0.15 points)
    # Check that statistic runs are bold, ~36pt, and white color
    try:
        auto_shapes = get_auto_shapes(slide3)
        format_matches = 0
        for shape in auto_shapes:
            text = get_shape_text(shape).strip()
            if text not in EXPECTED_STATS:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    is_bold = run.font.bold is True
                    # 36pt = 457200 EMU
                    is_36pt = run.font.size is not None and abs(run.font.size - Pt(36)) <= Pt(1)
                    is_white = False
                    try:
                        if run.font.color.type is not None:
                            rgb = str(run.font.color.rgb).upper()
                            is_white = rgb == 'FFFFFF'
                    except:
                        pass
                    if is_bold and is_36pt and is_white:
                        format_matches += 1
                    else:
                        print(f"  Format issue for '{text}': bold={is_bold}, 36pt={is_36pt}, white={is_white}")

        fmt_score = format_matches * (0.15 / 4.0)
        if format_matches >= 4:
            print(f"PASS: Component 6 — All 4 circles have correct text formatting (0.15 pts)")
            total_score += 0.15
        elif format_matches > 0:
            print(f"PARTIAL: Component 6 — {format_matches}/4 circles formatted correctly ({fmt_score:.4f} pts)")
            total_score += fmt_score
        else:
            print(f"FAIL: Component 6 — No circles have correct formatting (36pt bold white)")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
