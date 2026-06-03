"""
Initial Setup: Create UX Research presentation with 6 slides, slide 3 empty except title
Task ID: impress_stu_056
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_056'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines):
    """Add a slide with title and bullet-point body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = title_text
    body = slide.placeholders[1].text_frame
    body.text = body_lines[0]
    for line in body_lines[1:]:
        p = body.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "UX Research Report 2025"
    slide1.placeholders[1].text = "Customer Experience & Satisfaction Survey\nPrepared by the UX Research Team"

    # --- Slide 2: Research Methodology ---
    add_title_body_slide(prs, "Research Methodology", [
        "Online survey distributed to 1,200+ respondents across 5 regions",
        "Survey period: January 15 - February 28, 2025",
        "Mixed methods: quantitative Likert scales + qualitative open-ended questions",
        "Demographic segmentation by age group, region, and product usage frequency",
        "95% confidence interval with +/- 2.8% margin of error",
        "Follow-up interviews conducted with 45 participants for deeper insights",
    ])

    # --- Slide 3: Key Findings (EMPTY body, title only) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Findings"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: User Demographics ---
    add_title_body_slide(prs, "User Demographics", [
        "Age 18-24: 22% of respondents",
        "Age 25-34: 35% of respondents (largest segment)",
        "Age 35-44: 24% of respondents",
        "Age 45-54: 12% of respondents",
        "Age 55+: 7% of respondents",
        "Gender split: 52% female, 45% male, 3% non-binary/other",
    ])

    # --- Slide 5: Recommendations ---
    add_title_body_slide(prs, "Recommendations", [
        "Prioritize mobile-first redesign based on 68% mobile usage rate",
        "Implement personalized onboarding flow to address new user friction",
        "Expand self-service knowledge base (top request in qualitative feedback)",
        "Introduce live chat support during peak hours (10am - 2pm)",
        "Conduct quarterly pulse surveys to track satisfaction trends",
    ])

    # --- Slide 6: Thank You ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[0])
    slide6.shapes.title.text = "Thank You"
    slide6.placeholders[1].text = "Questions & Discussion\nux-research@company.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
