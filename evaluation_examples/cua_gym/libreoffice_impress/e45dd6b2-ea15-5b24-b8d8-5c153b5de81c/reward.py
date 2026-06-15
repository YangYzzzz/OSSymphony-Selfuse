"""
Reward Script: Apply underline and dark navy (#00008B) color to all text
               in the textbox and table on slide 5.
Task ID: osworld_impress_underline_darkred_table_009
Domain: libreoffice_impress
Scoring:
  Component 1: Textbox (TextBox 3) 4-paragraph text has underline=True      — 0.4 pts
  Component 2: Textbox (TextBox 3) 4-paragraph text has color=#00008B       — 0.3 pts
  Component 3: Table (3x3) all 9 cells have underline=True and color=#00008B — 0.3 pts
Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_009'
EXPECTED_COLOR = '00008B'
SLIDE_INDEX = 4  # Slide 5 (0-based)


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice state before verification."""
    try:
        import time
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion: all text in the content textbox and 3x3 table
    on slide 5 must have underline=True and color=#00008B.

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Expected at least 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]

    # Locate the content textbox (TextBox 3 — 4 bullet points) and the table
    content_textbox = None
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_shape = shape
        elif shape.has_text_frame and shape.shape_type == 17:  # TEXT_BOX
            # The content textbox has 4 paragraphs with black-text bullet points
            # (Not the "Key Metrics & Summary" header which has only 1 paragraph)
            if len(shape.text_frame.paragraphs) >= 4:
                content_textbox = shape

    if content_textbox is None:
        print("CRITICAL: Could not find content textbox with 4+ paragraphs on slide 5")
        print("REWARD: 0.0")
        return 0.0

    if table_shape is None:
        print("CRITICAL: Could not find table shape on slide 5")
        print("REWARD: 0.0")
        return 0.0

    # -----------------------------------------------------------------------
    # Component 1: Textbox paragraphs have underline=True on all runs (0.4 pts)
    # -----------------------------------------------------------------------
    try:
        paragraphs = content_textbox.text_frame.paragraphs
        tb_underline_count = 0
        tb_total_runs = 0
        tb_details = []
        for para in paragraphs:
            runs = [r for r in para.runs if (r.text or "").strip()]
            for run in runs:
                tb_total_runs += 1
                if run.font.underline is True:
                    tb_underline_count += 1
                    tb_details.append(f"PASS: '{run.text[:20]}' underline=True")
                else:
                    tb_details.append(f"FAIL: '{run.text[:20]}' underline={run.font.underline}")

        if tb_total_runs == 0:
            print("FAIL: Component 1 — No runs found in content textbox")
        elif tb_underline_count == tb_total_runs:
            print(f"PASS: Component 1 — All {tb_total_runs} textbox runs have underline=True (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — {tb_underline_count}/{tb_total_runs} textbox runs have underline=True")
            for d in tb_details:
                print(f"  {d}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: Textbox paragraphs have color=#00008B on all runs (0.3 pts)
    # -----------------------------------------------------------------------
    try:
        paragraphs = content_textbox.text_frame.paragraphs
        tb_color_count = 0
        tb_total_runs = 0
        tb_color_details = []
        for para in paragraphs:
            runs = [r for r in para.runs if (r.text or "").strip()]
            for run in runs:
                tb_total_runs += 1
                try:
                    if run.font.color.type is not None:
                        actual_color = str(run.font.color.rgb).upper()
                        if actual_color == EXPECTED_COLOR.upper():
                            tb_color_count += 1
                            tb_color_details.append(f"PASS: '{run.text[:20]}' color={actual_color}")
                        else:
                            tb_color_details.append(f"FAIL: '{run.text[:20]}' color={actual_color} (expected {EXPECTED_COLOR})")
                    else:
                        tb_color_details.append(f"FAIL: '{run.text[:20]}' color=inherited/None (expected {EXPECTED_COLOR})")
                except Exception as ce:
                    tb_color_details.append(f"FAIL: '{run.text[:20]}' color check error: {ce}")

        if tb_total_runs == 0:
            print("FAIL: Component 2 — No runs found in content textbox")
        elif tb_color_count == tb_total_runs:
            print(f"PASS: Component 2 — All {tb_total_runs} textbox runs have color=#{EXPECTED_COLOR} (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — {tb_color_count}/{tb_total_runs} textbox runs have color=#{EXPECTED_COLOR}")
            for d in tb_color_details:
                print(f"  {d}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Table all cells have underline=True AND color=#00008B (0.3 pts)
    # -----------------------------------------------------------------------
    try:
        table = table_shape.table
        num_rows = len(table.rows)
        num_cols = len(table.columns)
        tbl_pass_count = 0
        tbl_total_count = 0
        tbl_details = []

        for r in range(num_rows):
            for c in range(num_cols):
                cell = table.cell(r, c)
                for para in cell.text_frame.paragraphs:
                    runs = [run for run in para.runs if (run.text or "").strip()]
                    for run in runs:
                        tbl_total_count += 1
                        underline_ok = run.font.underline is True
                        color_ok = False
                        try:
                            if run.font.color.type is not None:
                                actual_color = str(run.font.color.rgb).upper()
                                color_ok = actual_color == EXPECTED_COLOR.upper()
                        except Exception:
                            pass

                        if underline_ok and color_ok:
                            tbl_pass_count += 1
                            tbl_details.append(f"PASS: cell({r},{c}) '{run.text[:15]}' underline+color OK")
                        else:
                            fail_reasons = []
                            if not underline_ok:
                                fail_reasons.append(f"underline={run.font.underline}")
                            if not color_ok:
                                try:
                                    actual_c = str(run.font.color.rgb).upper() if run.font.color.type is not None else "inherited"
                                    fail_reasons.append(f"color={actual_c}")
                                except Exception:
                                    fail_reasons.append("color=error")
                            tbl_details.append(f"FAIL: cell({r},{c}) '{run.text[:15]}' {'; '.join(fail_reasons)}")

        if tbl_total_count == 0:
            print("FAIL: Component 3 — No runs found in table cells")
        elif tbl_pass_count == tbl_total_count:
            print(f"PASS: Component 3 — All {tbl_total_count} table cell runs have underline+#{EXPECTED_COLOR} (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — {tbl_pass_count}/{tbl_total_count} table cell runs have both underline and color")
            for d in tbl_details:
                print(f"  {d}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on the VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
