"""
Initial Setup: 6-slide strategy deck for underline/dark navy formatting task
Task ID: osworld_impress_underline_darkred_table_009
Domain: libreoffice_impress

Creates a 6-slide strategy presentation with slide 5 containing:
- A content textbox with 4 bullet points (black text, no underline)
- A 3x3 summary table (black text, no underline)
Initial state has NO underline and NO #00008B color on slide 5 elements.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_underline_darkred_table_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_shape(shape, text, font_size=18, bold=False, color_rgb=None):
    """Helper to add text to a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.underline = False
    if color_rgb:
        run.font.color.rgb = color_rgb
    else:
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    HEADER_BLUE = RGBColor(0x1F, 0x3E, 0x6E)

    # ─────────────────────────────────────────────
    # Slide 1: Title Slide
    # ─────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Global Expansion Strategy 2025"
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = HEADER_BLUE
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    subtitle = slide1.placeholders[1]
    subtitle.text = "Prepared by the Strategic Planning Division\nQ1 2025"
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = DARK_GRAY

    # ─────────────────────────────────────────────
    # Slide 2: Executive Summary
    # ─────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = HEADER_BLUE
    content2 = slide2.placeholders[1]
    content2.text = "Our 2025 strategy focuses on three pillars"
    tf2 = content2.text_frame
    tf2.clear()
    points2 = [
        "Market expansion into Southeast Asia and Eastern Europe",
        "Digital transformation of core business processes",
        "Strategic partnerships with regional technology leaders",
        "Cost optimization through operational efficiency improvements",
    ]
    for i, pt in enumerate(points2):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = pt
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = BLACK

    # ─────────────────────────────────────────────
    # Slide 3: Market Analysis
    # ─────────────────────────────────────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Analysis"
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = HEADER_BLUE
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.clear()
    points3 = [
        "Total addressable market: $4.2B globally",
        "Southeast Asia projected growth: 18% CAGR through 2027",
        "Competitor landscape: 3 major players, 12 regional operators",
        "Regulatory environment: Favorable in target markets",
        "Customer acquisition cost trending down 12% YoY",
    ]
    for i, pt in enumerate(points3):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = pt
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = BLACK

    # ─────────────────────────────────────────────
    # Slide 4: Implementation Roadmap
    # ─────────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Implementation Roadmap"
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = HEADER_BLUE
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.clear()
    roadmap_items = [
        "Q1 2025: Infrastructure setup and team onboarding",
        "Q2 2025: Pilot launch in Singapore and Warsaw markets",
        "Q3 2025: Full regional rollout with localized offerings",
        "Q4 2025: Performance review and strategy refinement",
        "2026: Scaling operations and entering secondary markets",
    ]
    for i, item in enumerate(roadmap_items):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = item
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = BLACK

    # ─────────────────────────────────────────────
    # Slide 5: Key Metrics & Summary  <-- TASK TARGET
    # Content textbox with 4 points + 3x3 table
    # All text is BLACK, NO underline
    # ─────────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Title for slide 5
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.0), Inches(0.7))
    tf_title5 = title_box5.text_frame
    p_title5 = tf_title5.paragraphs[0]
    p_title5.text = "Key Metrics & Summary"
    run_title5 = p_title5.runs[0]
    run_title5.font.size = Pt(28)
    run_title5.font.bold = True
    run_title5.font.underline = False
    run_title5.font.color.rgb = HEADER_BLUE

    # Content textbox with 4 bullet points (black text, no underline)
    content_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5))
    tf5 = content_box5.text_frame
    tf5.word_wrap = True

    bullet_points = [
        "Revenue target: $125M by end of FY2025",
        "Customer base growth: 35% increase over prior year",
        "Operational efficiency: 20% cost reduction achieved",
        "Employee satisfaction score: 4.2 out of 5.0",
    ]

    for i, bp in enumerate(bullet_points):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        p.text = bp
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.bold = False
        run.font.underline = False
        run.font.color.rgb = BLACK  # Black, NOT dark navy

    # 3x3 Summary Table on slide 5 (black text, no underline)
    table_shape5 = slide5.shapes.add_table(
        3, 3,
        Inches(6.8), Inches(1.2),
        Inches(6.0), Inches(3.5)
    )
    table5 = table_shape5.table

    table_data = [
        ["Metric", "Target", "Actual"],
        ["Revenue Growth", "30%", "28%"],
        ["Market Share", "15%", "13%"],
    ]

    for r_idx, row_data in enumerate(table_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table5.cell(r_idx, c_idx)
            cell.text = cell_text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(15)
                    run.font.bold = (r_idx == 0)  # header row bold
                    run.font.underline = False
                    run.font.color.rgb = BLACK  # Black, NOT dark navy

    # ─────────────────────────────────────────────
    # Slide 6: Next Steps & Conclusion
    # ─────────────────────────────────────────────
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps & Conclusion"
    slide6.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = HEADER_BLUE
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.clear()
    next_steps = [
        "Finalize partnership agreements by March 31, 2025",
        "Complete regulatory filing for Singapore operations",
        "Launch recruitment drive for 45 regional positions",
        "Board presentation scheduled for April 15, 2025",
    ]
    for i, step in enumerate(next_steps):
        if i == 0:
            p = tf6.paragraphs[0]
        else:
            p = tf6.add_paragraph()
        p.text = step
        p.level = 1
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = BLACK

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
