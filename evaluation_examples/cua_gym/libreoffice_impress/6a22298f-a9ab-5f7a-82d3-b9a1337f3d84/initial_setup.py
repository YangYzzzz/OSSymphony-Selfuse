"""
Initial Setup: Board deck presentation with sequential slide numbering
Task ID: impress_fix_010
Domain: libreoffice_impress

Creates an 18-slide board deck with slide numbers visible on every slide,
numbered sequentially starting from 1.
"""

import os
import shlex
import subprocess
import time
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_slide_number_placeholder(slide, prs):
    """Add a slide number placeholder to a slide via XML using lxml."""
    import uuid as _uuid
    fld_id = str(_uuid.uuid4())
    off_x = int(prs.slide_width - Inches(1.5))
    off_y = int(prs.slide_height - Inches(0.5))
    cx = int(Inches(1.2))
    cy = int(Inches(0.35))
    sp_xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f'       xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f'       xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'  <p:nvSpPr>'
        f'    <p:cNvPr id="0" name="Slide Number Placeholder"/>'
        f'    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'    <p:nvPr><p:ph type="sldNum" sz="quarter" idx="12"/></p:nvPr>'
        f'  </p:nvSpPr>'
        f'  <p:spPr>'
        f'    <a:xfrm>'
        f'     <a:off x="{off_x}" y="{off_y}"/>'
        f'     <a:ext cx="{cx}" cy="{cy}"/>'
        f'    </a:xfrm>'
        f'  </p:spPr>'
        f'  <p:txBody>'
        f'    <a:bodyPr/>'
        f'    <a:lstStyle/>'
        f'    <a:p>'
        f'      <a:pPr algn="r"/>'
        f'      <a:fld id="{{{fld_id}}}" type="slidenum">'
        f'        <a:rPr lang="en-US" sz="1000" dirty="0"/>'
        f'        <a:t>&lt;#&gt;</a:t>'
        f'      </a:fld>'
        f'      <a:endParaRPr lang="en-US" sz="1000"/>'
        f'    </a:p>'
        f'  </p:txBody>'
        f'</p:sp>'
    )
    from lxml import etree as lxml_etree
    sp_elem = lxml_etree.fromstring(sp_xml)
    slide.shapes._spTree.append(sp_elem)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content definitions for an 18-slide board deck
    slide_data = [
        # Slide 1: Title
        {"layout": 5, "title": "Meridian Technologies Inc.", "subtitle": "Q1 2025 Board of Directors Meeting\nMarch 28, 2025\nConfidential"},
        # Slide 2: Agenda
        {"layout": 5, "title": "Agenda", "body": "1. Financial Performance Overview\n2. Product Development Update\n3. Market Expansion Strategy\n4. Talent & Organizational Health\n5. Risk Assessment & Compliance\n6. Technology Infrastructure\n7. Customer Success Metrics\n8. Strategic Partnerships\n9. Budget Allocation for Q2\n10. Sustainability Initiatives\n11. Competitive Landscape\n12. Board Action Items"},
        # Slide 3: Financial Overview
        {"layout": 5, "title": "Financial Performance Overview", "body": "Revenue: $47.3M (up 18% YoY)\nGross Margin: 68.2% (target: 65%)\nOperating Expenses: $28.1M\nNet Income: $8.4M\nCash Position: $124.5M\nBurn Rate: Positive cash flow achieved in February"},
        # Slide 4: Revenue Breakdown
        {"layout": 5, "title": "Revenue Breakdown by Segment", "body": "Enterprise SaaS: $28.7M (61%)\nMid-Market: $11.2M (24%)\nSMB Self-Serve: $5.8M (12%)\nProfessional Services: $1.6M (3%)\n\nEnterprise deal count: 47 new contracts\nAverage contract value: $610K (up from $485K)"},
        # Slide 5: Product Update
        {"layout": 5, "title": "Product Development Update", "body": "Platform v4.2 launched February 12\n- AI-powered analytics dashboard\n- Real-time collaboration features\n- SOC 2 Type II compliance module\n\nUpcoming: v4.3 targeted for May release\n- Predictive forecasting engine\n- Advanced API gateway\n- Mobile app redesign"},
        # Slide 6: Market Expansion
        {"layout": 5, "title": "Market Expansion Strategy", "body": "EMEA Region:\n- London office opened January 15\n- 12 enterprise prospects in pipeline\n- Hired VP of EMEA Sales: Patricia Okafor\n\nAPAC Region:\n- Singapore partnership signed with TechBridge Asia\n- Japan market entry planned for Q3\n- Regulatory approvals pending in South Korea"},
        # Slide 7: Talent
        {"layout": 5, "title": "Talent & Organizational Health", "body": "Headcount: 312 (up from 278 in Q4)\nNew Hires: 41 in Q1\nAttrition Rate: 8.2% annualized\nEmployee NPS: 72 (industry avg: 45)\n\nKey Hires:\n- CTO: Dr. Amir Rashidi (ex-Google DeepMind)\n- VP Engineering: Lisa Yamamoto\n- Head of Data Science: Carlos Mendez"},
        # Slide 8: Risk Assessment
        {"layout": 5, "title": "Risk Assessment & Compliance", "body": "Cybersecurity: Zero breaches in Q1\nGDPR Compliance: Full audit passed\nSOX Readiness: 94% controls documented\n\nEmerging Risks:\n- Supply chain dependency on 3 cloud providers\n- Regulatory changes in EU AI Act\n- Currency exposure from EMEA expansion\n\nMitigation: Multi-cloud strategy approved, legal review underway"},
        # Slide 9: Technology Infrastructure
        {"layout": 5, "title": "Technology Infrastructure", "body": "Uptime: 99.97% (SLA target: 99.95%)\nLatency: P99 at 142ms (down from 198ms)\nCloud Spend: $2.1M/month\n\nInfrastructure Initiatives:\n- Migration to Kubernetes completed\n- Edge computing nodes deployed in 8 regions\n- Disaster recovery tested successfully on Feb 20"},
        # Slide 10: Customer Success
        {"layout": 5, "title": "Customer Success Metrics", "body": "Net Revenue Retention: 118%\nCustomer Satisfaction (CSAT): 4.6/5.0\nSupport Ticket Resolution: Avg 4.2 hours\nChurn Rate: 2.1% quarterly\n\nTop Account Renewals:\n- GlobalBank Corp: $3.2M (3-year renewal)\n- Nexus Healthcare: $1.8M (expanded scope)\n- Pinnacle Retail Group: $1.1M (new modules added)"},
        # Slide 11: Strategic Partnerships
        {"layout": 5, "title": "Strategic Partnerships", "body": "Active Partnerships: 14\n\nNew in Q1:\n- Microsoft Azure Marketplace integration\n- Salesforce AppExchange listing\n- Deloitte implementation partnership\n\nPipeline Impact: $12.4M influenced revenue\nCo-marketing campaigns: 6 joint webinars, 2 whitepapers"},
        # Slide 12: Budget Q2
        {"layout": 5, "title": "Q2 2025 Budget Allocation", "body": "Total Budget: $31.5M\n\nEngineering: $12.8M (41%)\nSales & Marketing: $9.2M (29%)\nG&A: $4.7M (15%)\nCustomer Success: $3.1M (10%)\nR&D Innovation Fund: $1.7M (5%)\n\nCapital Expenditure: $2.4M (data center expansion)"},
        # Slide 13: Sustainability
        {"layout": 5, "title": "Sustainability Initiatives", "body": "Carbon Footprint: Reduced 22% vs Q1 2024\nRenewable Energy: 78% of operations\nPaper-free Office: Achieved across all locations\n\n2025 Goals:\n- Carbon neutral by December 2025\n- 100% renewable energy by Q3\n- ESG report publication in June\n- Community investment: $500K allocated"},
        # Slide 14: Competitive Landscape
        {"layout": 5, "title": "Competitive Landscape Analysis", "body": "Market Position: #3 in enterprise segment\n\nCompetitor Updates:\n- Apex Solutions: Raised $200M Series D\n- DataForge: Acquired CloudSync for $450M\n- Synergy Platform: Leadership changes, 3 exec departures\n\nOur Differentiators:\n- AI-first architecture\n- Superior time-to-value (avg 6 weeks vs 12 weeks)\n- Highest customer satisfaction in segment"},
        # Slide 15: Key Metrics Dashboard
        {"layout": 5, "title": "Key Metrics Dashboard", "body": "ARR: $189.2M\nMRR Growth: 4.2% month-over-month\nCAC Payback: 14 months\nLTV/CAC Ratio: 5.8x\nRule of 40: Score 52\nDAU/MAU Ratio: 0.62\n\nAll metrics trending above board-approved targets"},
        # Slide 16: Investor Relations
        {"layout": 5, "title": "Investor Relations Update", "body": "Secondary Share Sale: Completed Feb 15\n- $18M in liquidity for early employees\n- Valuation: $1.2B (post-money)\n\nIPO Readiness:\n- S-1 draft in progress with Morgan Stanley\n- Target filing: Q4 2025\n- Audit committee review scheduled for April\n\nInvestor Communications: 8 board updates sent in Q1"},
        # Slide 17: Action Items
        {"layout": 5, "title": "Board Action Items", "body": "1. VOTE: Approve Q2 budget allocation ($31.5M)\n2. VOTE: Authorize EMEA office lease extension (5-year term)\n3. REVIEW: IPO timeline and underwriter selection\n4. DISCUSS: AI governance framework proposal\n5. APPROVE: Executive compensation adjustments\n6. NOTE: Next board meeting - June 27, 2025"},
        # Slide 18: Closing
        {"layout": 5, "title": "Thank You", "body": "Questions & Discussion\n\nContact: board-relations@meridiantech.com\nBoard Portal: https://board.meridiantech.com\n\nNext Meeting: June 27, 2025\nLocation: San Francisco HQ / Virtual Option Available"},
    ]

    for i, sd in enumerate(slide_data):
        slide = prs.slides.add_slide(prs.slide_layouts[sd["layout"]])

        # Title text box
        title_top = Inches(0.5) if i > 0 else Inches(2.0)
        title_size = Pt(36) if i == 0 else Pt(28)
        txBox = slide.shapes.add_textbox(Inches(0.8), title_top, Inches(11.5), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd["title"]
        p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = title_size
        run.font.bold = True
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

        # Body or subtitle
        body_text = sd.get("body") or sd.get("subtitle")
        if body_text:
            body_top = Inches(1.8) if i > 0 else Inches(3.5)
            bBox = slide.shapes.add_textbox(Inches(0.8), body_top, Inches(11.5), Inches(5.0))
            btf = bBox.text_frame
            btf.word_wrap = True
            for li, line in enumerate(body_text.split('\n')):
                if li == 0:
                    bp = btf.paragraphs[0]
                else:
                    bp = btf.add_paragraph()
                bp.text = line
                bp.space_after = Pt(4)
                if line.strip():
                    br = bp.runs[0]
                    br.font.size = Pt(16) if i > 0 else Pt(20)
                    br.font.name = "Calibri"
                    br.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        # Add slide number placeholder to every slide
        add_slide_number_placeholder(slide, prs)

    prs.save(OUTPUT)

    print(f'Initial file created: {OUTPUT}')
    print(f'Slides: 18, all with slide number placeholders')

    # Launch GUI
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
