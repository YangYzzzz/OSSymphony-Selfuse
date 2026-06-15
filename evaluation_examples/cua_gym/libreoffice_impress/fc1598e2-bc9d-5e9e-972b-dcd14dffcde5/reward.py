"""
Reward Script: Organizational chart on slide 3 of New_Hire_Onboarding.pptx
Task ID: impress_ps_022
Domain: libreoffice_impress
Scoring:
  Component 1: CEO box present with correct text (0.15)
  Component 2: 3 VP boxes with correct titles (0.20)
  Component 3: 6 Director boxes with correct titles (0.25)
  Component 4: All boxes are rounded rectangles with 12pt font (0.10)
  Component 5: 9 line connectors on slide 3 (0.20)
  Component 6: Hierarchical vertical layout (CEO top, VPs mid, Dirs bottom) (0.10)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_022'


def persist_app_state(domain):
    """Try to save any unsaved GUI state via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print("PERSIST: ctrl+s sent for %s" % domain)
        except Exception as e:
            print("PERSIST_WARN: save hook failed: %s" % e)


def verify_task(file_path):
    """
    Verify organizational chart on slide 3.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print("CRITICAL: Cannot import python-pptx: %s" % e)
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print("FAIL: Presentation has fewer than 3 slides (%d)" % len(prs.slides))
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # 0-indexed, slide 3

    # Collect all auto shapes (boxes) and lines (connectors) on slide 3
    auto_shapes = []
    lines = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            auto_shapes.append({
                'text': text,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'name': shape.name,
                'shape': shape,
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            lines.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
            })

    print("Found %d auto_shapes and %d lines on slide 3" % (len(auto_shapes), len(lines)))

    # Build text lookup (case-insensitive for matching)
    shape_texts = {s['text'].lower(): s for s in auto_shapes if s['text']}

    # -------------------------------------------------------
    # Component 1: CEO box present with correct text (0.15)
    # -------------------------------------------------------
    try:
        ceo_found = False
        for s in auto_shapes:
            if 'ceo' in s['text'].lower():
                ceo_found = True
                print("PASS: Component 1 — CEO box found with text '%s' (0.15 pts)" % s['text'])
                total_score += 0.15
                break
        if not ceo_found:
            print("FAIL: Component 1 — No box with 'CEO' text found on slide 3")
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    # -------------------------------------------------------
    # Component 2: 3 VP boxes with correct titles (0.20)
    # -------------------------------------------------------
    try:
        expected_vps = ['vp engineering', 'vp sales', 'vp operations']
        found_vps = []
        for s in auto_shapes:
            txt = s['text'].lower()
            for vp in expected_vps:
                if vp in txt:
                    found_vps.append(s['text'])
                    break
        vp_count = len(found_vps)
        if vp_count == 3:
            print("PASS: Component 2 — All 3 VP boxes found: %s (0.20 pts)" % found_vps)
            total_score += 0.20
        elif vp_count > 0:
            partial = round(0.20 * vp_count / 3, 2)
            print("PARTIAL: Component 2 — %d/3 VP boxes found: %s (%s pts)" % (vp_count, found_vps, partial))
            total_score += partial
        else:
            print("FAIL: Component 2 — No VP boxes found")
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    # -------------------------------------------------------
    # Component 3: 6 Director boxes with correct titles (0.25)
    # -------------------------------------------------------
    try:
        director_boxes = []
        for s in auto_shapes:
            txt = s['text'].lower()
            if 'dir' in txt and 'ceo' not in txt and 'vp' not in txt:
                director_boxes.append(s['text'])
        dir_count = len(director_boxes)
        if dir_count >= 6:
            print("PASS: Component 3 — %d director boxes found: %s (0.25 pts)" % (dir_count, director_boxes))
            total_score += 0.25
        elif dir_count > 0:
            partial = round(0.25 * min(dir_count, 6) / 6, 2)
            print("PARTIAL: Component 3 — %d/6 director boxes found: %s (%s pts)" % (dir_count, director_boxes, partial))
            total_score += partial
        else:
            print("FAIL: Component 3 — No director boxes found on slide 3")
    except Exception as e:
        print("ERROR: Component 3 — %s" % e)

    # -------------------------------------------------------
    # Component 4: All boxes use 12pt font (0.10)
    # -------------------------------------------------------
    try:
        total_boxes_with_text = 0
        boxes_with_12pt = 0
        for s in auto_shapes:
            if not s['text']:
                continue
            total_boxes_with_text += 1
            shape_obj = s['shape']
            if shape_obj.has_text_frame:
                for para in shape_obj.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            pt_size = round(run.font.size / 12700, 1)
                            if abs(pt_size - 12.0) < 0.5:
                                boxes_with_12pt += 1
                            break
                    break  # only check first run of first paragraph

        if total_boxes_with_text > 0 and boxes_with_12pt == total_boxes_with_text:
            print("PASS: Component 4 — All %d text boxes use 12pt font (0.10 pts)" % total_boxes_with_text)
            total_score += 0.10
        elif boxes_with_12pt > 0:
            partial = round(0.10 * boxes_with_12pt / max(total_boxes_with_text, 1), 2)
            print("PARTIAL: Component 4 — %d/%d boxes have 12pt font (%s pts)" % (boxes_with_12pt, total_boxes_with_text, partial))
            total_score += partial
        else:
            print("FAIL: Component 4 — No boxes with 12pt font found (checked %d boxes)" % total_boxes_with_text)
    except Exception as e:
        print("ERROR: Component 4 — %s" % e)

    # -------------------------------------------------------
    # Component 5: 9 line connectors on slide 3 (0.20)
    # -------------------------------------------------------
    try:
        line_count = len(lines)
        if line_count >= 9:
            print("PASS: Component 5 — %d line connectors found (0.20 pts)" % line_count)
            total_score += 0.20
        elif line_count > 0:
            partial = round(0.20 * min(line_count, 9) / 9, 2)
            print("PARTIAL: Component 5 — %d/9 line connectors found (%s pts)" % (line_count, partial))
            total_score += partial
        else:
            print("FAIL: Component 5 — No line connectors found on slide 3")
    except Exception as e:
        print("ERROR: Component 5 — %s" % e)

    # -------------------------------------------------------
    # Component 6: Hierarchical vertical layout (0.10)
    # CEO at top, VPs in middle, Directors at bottom
    # -------------------------------------------------------
    try:
        ceo_top = None
        vp_tops = []
        dir_tops = []
        for s in auto_shapes:
            txt = s['text'].lower()
            if 'ceo' in txt:
                ceo_top = s['top']
            elif 'vp' in txt:
                vp_tops.append(s['top'])
            elif 'dir' in txt:
                dir_tops.append(s['top'])

        hierarchy_ok = False
        if ceo_top is not None and len(vp_tops) >= 3 and len(dir_tops) >= 6:
            avg_vp_top = sum(vp_tops) / len(vp_tops)
            avg_dir_top = sum(dir_tops) / len(dir_tops)
            # CEO should be above VPs, VPs should be above Directors
            if ceo_top < avg_vp_top and avg_vp_top < avg_dir_top:
                hierarchy_ok = True

        if hierarchy_ok:
            print("PASS: Component 6 — Hierarchical layout verified (CEO top=%d, VP avg=%d, Dir avg=%d) (0.10 pts)" %
                  (ceo_top, avg_vp_top, avg_dir_top))
            total_score += 0.10
        else:
            print("FAIL: Component 6 — Hierarchy not verified. CEO=%s, VPs=%d tops, Dirs=%d tops" %
                  (ceo_top, len(vp_tops), len(dir_tops)))
    except Exception as e:
        print("ERROR: Component 6 — %s" % e)

    final_score = round(min(total_score, 1.0), 2)
    print()
    print("Score: %s/1.0" % final_score)
    print("REWARD: %s" % final_score)
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = '%s/%s.pptx' % (WORKDIR, TASK_ID)
if not os.path.exists(file_path):
    # Try alternate names
    alt = '%s/New_Hire_Onboarding.pptx' % WORKDIR
    if os.path.exists(alt):
        file_path = alt
    else:
        print("File not found: %s" % file_path)
        print("REWARD: 0.0")
        exit()

verify_task(file_path)
