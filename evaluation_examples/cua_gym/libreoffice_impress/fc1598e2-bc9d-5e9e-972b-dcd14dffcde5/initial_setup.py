"""
Initial Setup: Build a 10-slide New Hire Onboarding presentation.
Slide 3 has title 'Our Organization' with empty content area.
Task ID: impress_ps_022
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_022'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Slide with title only and empty content area."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add title textbox manually
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Welcome to Apex Dynamics",
                    "New Hire Onboarding Program — Q2 2025")

    # Slide 2: Agenda
    add_content_slide(prs, "Onboarding Agenda", [
        "Company Overview & Mission",
        "Our Organization",
        "Benefits & Compensation",
        "IT Setup & Security Policies",
        "Team Introductions",
        "30-60-90 Day Plan",
        "Resources & Contacts",
        "Q&A Session",
    ])

    # Slide 3: "Our Organization" — title only, empty content area (NO org chart)
    add_title_only_slide(prs, "Our Organization")

    # Slide 4: Company History
    add_content_slide(prs, "Company History", [
        "Founded in 2008 by Dr. Elena Vasquez in Austin, TX",
        "Expanded to 3 global offices by 2015 (London, Singapore, Austin)",
        "Reached $500M annual revenue in 2022",
        "Currently 2,400+ employees across 12 countries",
        "Named 'Best Workplace in Tech' by Forbes 2023 & 2024",
    ])

    # Slide 5: Benefits Overview
    add_content_slide(prs, "Benefits & Compensation", [
        "Comprehensive health, dental, and vision insurance",
        "401(k) with 6% company match (vested immediately)",
        "Unlimited PTO policy with 15-day minimum encouraged",
        "Annual learning & development budget: $3,000 per employee",
        "Home office stipend: $1,500 for new hires",
        "Employee stock purchase plan (15% discount)",
    ])

    # Slide 6: IT Setup
    add_content_slide(prs, "IT Setup & Security", [
        "Laptop will be provisioned within 24 hours of start date",
        "Use Okta SSO for all company applications",
        "VPN required for accessing internal tools remotely",
        "Mandatory security awareness training (due within 7 days)",
        "Report suspicious emails to security@apexdynamics.com",
        "Two-factor authentication required on all accounts",
    ])

    # Slide 7: Team Culture
    add_content_slide(prs, "Our Culture & Values", [
        "Innovation First — we encourage bold ideas and experiments",
        "Customer Obsession — every decision starts with the customer",
        "Radical Transparency — open communication at every level",
        "Continuous Growth — invest in yourself and your teammates",
        "Collaboration over competition — we win as a team",
    ])

    # Slide 8: 30-60-90 Day Plan
    add_content_slide(prs, "Your 30-60-90 Day Plan", [
        "Days 1-30: Complete onboarding, meet your team, shadow projects",
        "Days 31-60: Take ownership of first deliverable, attend cross-team meetings",
        "Days 61-90: Lead a small initiative, present learnings to your manager",
        "Quarterly check-ins with HR and your direct manager",
        "Annual performance review aligned to company OKRs",
    ])

    # Slide 9: Resources
    add_content_slide(prs, "Resources & Key Contacts", [
        "HR Portal: hr.apexdynamics.com",
        "IT Help Desk: helpdesk@apexdynamics.com | Ext. 4500",
        "Your HR Business Partner: Priya Nakamura (priya.n@apexdynamics.com)",
        "Facilities Manager: Tom Eriksson (tom.e@apexdynamics.com)",
        "Employee Handbook: Available on the HR Portal under 'Documents'",
    ])

    # Slide 10: Q&A / Closing
    add_title_slide(prs, "Questions?",
                    "Thank you for joining Apex Dynamics!\nWe're excited to have you on the team.")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
