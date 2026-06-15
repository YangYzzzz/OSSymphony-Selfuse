"""
Initial Setup: Academic conference talk presentation (7 slides)
Task ID: osworld_impress_multi_op_combined_010
Domain: libreoffice_impress

Creates a 7-slide academic conference talk presentation.
- Slides 2 and 3: plain black titles (NOT bold, NOT green)
- Slide 4: title 'Discussion' with NO speaker notes
- Slide 6: body textbox reading 'TODO: update'
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, title_text, content_lines):
    """Add a slide with title and bulleted content."""
    layout = prs.slide_layouts[1]  # Title + Content layout
    slide = prs.slides.add_slide(layout)
    title_shape = slide.shapes.title
    title_shape.text = title_text
    # Plain black title (no bold, no color override)
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = content_lines[0] if content_lines else ""
    for line in content_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    layout0 = prs.slide_layouts[0]  # Title Slide
    slide1 = prs.slides.add_slide(layout0)
    slide1.shapes.title.text = "Advances in Distributed Machine Learning"
    slide1.placeholders[1].text = "Dr. Priya Nambiar, Dr. Leon Fischer\nInternational Conference on AI Systems, 2025"

    # --- Slide 2: Introduction (plain black title — NOT bold, NOT green) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    for para in slide2.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    body2 = slide2.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = "Motivation: scaling ML to thousands of nodes"
    p = tf2.add_paragraph()
    p.text = "Prior work: data parallelism, model parallelism"
    p = tf2.add_paragraph()
    p.text = "Our contribution: adaptive gradient compression"
    p = tf2.add_paragraph()
    p.text = "Evaluation on ImageNet and BERT fine-tuning tasks"

    # --- Slide 3: Related Work (plain black title — NOT bold, NOT green) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Related Work"
    for para in slide3.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.bold = False
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "Horovod (Sergeev et al., 2018): ring-allreduce"
    p = tf3.add_paragraph()
    p.text = "PowerSGD (Vogels et al., 2019): low-rank approximation"
    p = tf3.add_paragraph()
    p.text = "Top-K sparsification (Alistarh et al., 2017)"
    p = tf3.add_paragraph()
    p.text = "Gradient checkpointing: memory vs. compute tradeoff"

    # --- Slide 4: Discussion (title = 'Discussion', NO speaker notes) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Discussion"
    body4 = slide4.placeholders[1]
    tf4 = body4.text_frame
    tf4.text = "Limitations of current approach"
    p = tf4.add_paragraph()
    p.text = "Bandwidth vs. accuracy tradeoff at high compression ratios"
    p = tf4.add_paragraph()
    p.text = "Sensitivity to learning rate scheduling"
    p = tf4.add_paragraph()
    p.text = "Future directions: federated learning extensions"
    # NOTE: intentionally NO speaker notes on slide 4

    # --- Slide 5: Methodology ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Methodology"
    body5 = slide5.placeholders[1]
    tf5 = body5.text_frame
    tf5.text = "Adaptive threshold τ based on gradient norm"
    p = tf5.add_paragraph()
    p.text = "Error feedback mechanism to preserve convergence"
    p = tf5.add_paragraph()
    p.text = "Warm-up phase: full gradients for first 5 epochs"
    p = tf5.add_paragraph()
    p.text = "Asynchronous momentum correction"

    # --- Slide 6: Experimental Results (with 'TODO: update' textbox) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Experimental Results"
    body6 = slide6.placeholders[1]
    tf6 = body6.text_frame
    tf6.text = "Benchmark: 64-node GPU cluster (Nvidia A100)"
    p = tf6.add_paragraph()
    p.text = "Top-1 accuracy: 76.3% on ImageNet (baseline: 76.1%)"
    p = tf6.add_paragraph()
    p.text = "Training speedup: 2.4x vs. Horovod at 95% sparsity"
    p = tf6.add_paragraph()
    p.text = "Communication overhead reduced by 18%"

    # Add the 'TODO: update' text box on slide 6
    todo_left = Inches(1.0)
    todo_top = Inches(6.0)
    todo_width = Inches(4.0)
    todo_height = Inches(0.6)
    todo_txbox = slide6.shapes.add_textbox(todo_left, todo_top, todo_width, todo_height)
    todo_tf = todo_txbox.text_frame
    todo_tf.text = "TODO: update"

    # --- Slide 7: Conclusion ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Conclusion"
    body7 = slide7.placeholders[1]
    tf7 = body7.text_frame
    tf7.text = "Adaptive gradient compression improves distributed training efficiency"
    p = tf7.add_paragraph()
    p.text = "No significant accuracy degradation observed"
    p = tf7.add_paragraph()
    p.text = "Code and datasets available at github.com/dml-lab/agc"
    p = tf7.add_paragraph()
    p.text = "Acknowledgements: NSF Grant #1234567, compute time from Cloud Labs"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup — open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
