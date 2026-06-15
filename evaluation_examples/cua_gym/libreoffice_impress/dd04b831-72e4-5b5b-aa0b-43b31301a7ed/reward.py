"""
Reward Script: Format titles on slides 2 and 3 with bold + green color,
               add slide title as speaker note on slide 4,
               and delete 'TODO: update' textbox from slide 6.
Task ID: osworld_impress_multi_op_combined_010
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 title is bold AND green (0.30 points)
  Component 2: Slide 3 title is bold AND green (0.30 points)
  Component 3: Slide 4 speaker notes contain 'Discussion' (0.20 points)
  Component 4: 'TODO: update' textbox removed from slide 6 (0.20 points)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_multi_op_combined_010'

# Expected green color from golden env
EXPECTED_GREEN = '008000'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 title is bold AND green (0.30 points)
    # This FAILS on initial (bold=False, color=000000) and PASSES on golden (bold=True, color=008000)
    try:
        slide2 = prs.slides[1]
        title2 = None
        for shape in slide2.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                # Check if it's the title placeholder (idx 0)
                try:
                    ph = shape.placeholder_format
                    if ph.idx == 0:
                        title2 = shape
                        break
                except Exception:
                    pass

        if title2 is None:
            print("FAIL: Component 1 — Could not find title placeholder on slide 2")
        else:
            # Check runs for bold and green color
            all_bold = True
            all_green = True
            run_count = 0
            for para in title2.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    run_count += 1
                    # Check bold: None is treated as "not bold"
                    is_bold = run.font.bold is True
                    if not is_bold:
                        all_bold = False
                    # Check green color
                    try:
                        color_type = run.font.color.type
                        if color_type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str != EXPECTED_GREEN.upper():
                                all_green = False
                        else:
                            all_green = False
                    except Exception:
                        all_green = False

            if run_count == 0:
                print("FAIL: Component 1 — No text runs found in slide 2 title")
            elif all_bold and all_green:
                print(f"PASS: Component 1 — Slide 2 title is bold and green (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Slide 2 title: bold={all_bold}, green={all_green}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 title is bold AND green (0.30 points)
    # This FAILS on initial (bold=False, color=000000) and PASSES on golden (bold=True, color=008000)
    try:
        slide3 = prs.slides[2]
        title3 = None
        for shape in slide3.shapes:
            if shape.is_placeholder and shape.has_text_frame:
                try:
                    ph = shape.placeholder_format
                    if ph.idx == 0:
                        title3 = shape
                        break
                except Exception:
                    pass

        if title3 is None:
            print("FAIL: Component 2 — Could not find title placeholder on slide 3")
        else:
            all_bold = True
            all_green = True
            run_count = 0
            for para in title3.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    run_count += 1
                    is_bold = run.font.bold is True
                    if not is_bold:
                        all_bold = False
                    try:
                        color_type = run.font.color.type
                        if color_type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str != EXPECTED_GREEN.upper():
                                all_green = False
                        else:
                            all_green = False
                    except Exception:
                        all_green = False

            if run_count == 0:
                print("FAIL: Component 2 — No text runs found in slide 3 title")
            elif all_bold and all_green:
                print(f"PASS: Component 2 — Slide 3 title is bold and green (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 2 — Slide 3 title: bold={all_bold}, green={all_green}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 4 speaker notes contain 'Discussion' (0.20 points)
    # This FAILS on initial (empty notes) and PASSES on golden (notes = 'Discussion')
    try:
        slide4 = prs.slides[3]
        notes_text = ""
        try:
            notes_text = slide4.notes_slide.notes_text_frame.text.strip()
        except Exception as ne:
            print(f"WARN: Component 3 — Could not read notes: {ne}")

        if "Discussion" in notes_text:
            print(f"PASS: Component 3 — Slide 4 notes contain 'Discussion' (value: {repr(notes_text)}) (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Slide 4 notes: expected 'Discussion', found: {repr(notes_text)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 'TODO: update' textbox removed from slide 6 (0.20 points)
    # This FAILS on initial (TextBox 3 with 'TODO: update' present) and PASSES on golden (removed)
    try:
        slide6 = prs.slides[5]
        todo_found = False
        for shape in slide6.shapes:
            if not shape.is_placeholder and shape.has_text_frame:
                # Non-placeholder textbox — check if it contains 'TODO: update'
                full_text = shape.text_frame.text.strip()
                if "TODO" in full_text and "update" in full_text.lower():
                    todo_found = True
                    print(f"FAIL: Component 4 — 'TODO: update' textbox still present on slide 6 (shape: {shape.name}, text: {repr(full_text)})")
                    break

        if not todo_found:
            print(f"PASS: Component 4 — 'TODO: update' textbox is absent from slide 6 (0.20 pts)")
            total_score += 0.20
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
