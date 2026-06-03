"""
Initial Setup: Center table on slide 3
Task ID: impress_tct_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard slide size: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Workshop Schedule"
    slide1.placeholders[1].text = "Annual Professional Development Series 2025"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "This year we are offering five hands-on workshops covering data analytics, leadership, cloud computing, UX design, and cybersecurity."
    p2 = body2.add_paragraph()
    p2.text = "Each session runs for a full day and includes lunch and materials."
    p3 = body2.add_paragraph()
    p3.text = "Registration is open to all departments on a first-come, first-served basis."

    # --- Slide 3: Schedule Table (top-left, NOT centered) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Detailed Schedule"
    run = p.runs[0]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Table: 5 rows x 3 columns, positioned at top-left corner
    rows, cols = 5, 3
    table_width = Inches(7.0)
    table_height = Inches(3.0)
    table_left = Inches(0.5)
    table_top = Inches(0.5)

    table_shape = slide3.shapes.add_table(
        rows, cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.0)

    # Headers
    headers = ["Workshop", "Instructor", "Date"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Header background
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Data rows
    data = [
        ["Data Analytics Fundamentals", "Dr. Sarah Chen", "March 15, 2025"],
        ["Leadership & Communication", "Marcus Johnson", "April 22, 2025"],
        ["Cloud Computing Essentials", "Priya Patel", "May 10, 2025"],
        ["UX Design Principles", "Elena Rodriguez", "June 18, 2025"],
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Alternating row colors
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)

    # --- Slide 4: Closing ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Registration & Contact"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Register through the company intranet portal by selecting Professional Development."
    p4 = body4.add_paragraph()
    p4.text = "For questions, contact training@company.com or ext. 4521."
    p5 = body4.add_paragraph()
    p5.text = "Spots are limited to 30 participants per workshop."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
