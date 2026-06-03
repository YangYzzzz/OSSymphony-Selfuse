"""
Reward Script: Center the table horizontally and vertically on slide 3
Task ID: impress_tct_010
Domain: libreoffice_impress
Scoring:
  Component 1: Table horizontally centered on slide 3 (0.5 pts)
  Component 2: Table vertically centered on slide 3 (0.5 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_010'


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # 0-indexed, slide 3

    # Find the table shape on slide 3
    table_shape = None
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table_shape = shape
            break

    if table_shape is None:
        print("FAIL: No table found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    table_left = table_shape.left
    table_top = table_shape.top
    table_width = table_shape.width
    table_height = table_shape.height

    # Expected centered positions
    expected_left = (slide_width - table_width) // 2
    expected_top = (slide_height - table_height) // 2

    print(f"Slide dimensions: {slide_width} x {slide_height} EMU ({slide_width/914400:.2f} x {slide_height/914400:.2f} in)")
    print(f"Table dimensions: {table_width} x {table_height} EMU ({table_width/914400:.2f} x {table_height/914400:.2f} in)")
    print(f"Table position: left={table_left} EMU ({table_left/914400:.4f} in), top={table_top} EMU ({table_top/914400:.4f} in)")
    print(f"Expected center: left={expected_left} EMU ({expected_left/914400:.4f} in), top={expected_top} EMU ({expected_top/914400:.4f} in)")

    # Tolerance: 5% of slide dimension (generous for GUI centering)
    h_tolerance = int(slide_width * 0.05)
    v_tolerance = int(slide_height * 0.05)

    # Component 1: Table is horizontally centered (0.5 points)
    try:
        h_diff = abs(table_left - expected_left)
        print(f"Horizontal offset from center: {h_diff} EMU ({h_diff/914400:.4f} in)")
        if h_diff <= h_tolerance:
            print(f"PASS: Component 1 -- Table horizontally centered (within {h_tolerance} EMU tolerance) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 -- Table NOT horizontally centered. Off by {h_diff} EMU ({h_diff/914400:.4f} in), tolerance={h_tolerance} EMU")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Table is vertically centered (0.5 points)
    try:
        v_diff = abs(table_top - expected_top)
        print(f"Vertical offset from center: {v_diff} EMU ({v_diff/914400:.4f} in)")
        if v_diff <= v_tolerance:
            print(f"PASS: Component 2 -- Table vertically centered (within {v_tolerance} EMU tolerance) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 -- Table NOT vertically centered. Off by {v_diff} EMU ({v_diff/914400:.4f} in), tolerance={v_tolerance} EMU")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
