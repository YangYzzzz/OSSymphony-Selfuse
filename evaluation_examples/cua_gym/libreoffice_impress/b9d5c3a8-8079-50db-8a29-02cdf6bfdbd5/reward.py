"""
Reward Script: Statistics highlight on slide 7 with 97%, underline, and subtitle
Task ID: impress_design_085
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): '97%' text with 96pt bold #2C3E50 centered
  Component 2 (0.35): Thin rectangle underline ~8in wide, #E74C3C fill, at y=4in
  Component 3 (0.30): 'Client Satisfaction Rate' in 24pt #666666 centered at y=4.5in
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_design_085'


def persist_app_state(domain: str):
    """Save any open LibreOffice document before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def is_approx_equal(val1, val2, tolerance=0.05):
    """Check if two numeric values are approximately equal (5% tolerance)."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 50000  # ~0.05 inches in EMU
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # 0-indexed, slide 7
    shapes = list(slide.shapes)

    # Helper: find text shapes containing specific text
    def find_text_shapes(target_text):
        found = []
        for shape in shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if target_text.lower() in para.text.strip().lower():
                        found.append(shape)
                        break
        return found

    # Component 1: '97%' text with correct formatting (0.35 points)
    try:
        matches_97 = find_text_shapes('97%')
        if not matches_97:
            print("FAIL: Component 1 -- No shape with '97%' text found on slide 7")
        else:
            shape_97 = matches_97[0]
            para = None
            for p in shape_97.text_frame.paragraphs:
                if '97%' in p.text.strip():
                    para = p
                    break

            if para is None:
                print("FAIL: Component 1 -- Could not find paragraph with '97%'")
            else:
                comp1_score = 0.0
                comp1_max = 0.35

                # Check text content
                runs = [r for r in para.runs if r.text.strip()]
                full_text = ''.join(r.text for r in runs).strip()
                if '97%' in full_text:
                    comp1_score += 0.05
                    print(f"PASS: Component 1a -- Text contains '97%' (found: '{full_text}')")
                else:
                    print(f"FAIL: Component 1a -- Expected '97%', found '{full_text}'")

                # Check font size (96pt = 1219200 EMU)
                if runs:
                    font_size = runs[0].font.size
                    if font_size is not None and is_approx_equal(font_size, 1219200):
                        comp1_score += 0.07
                        print(f"PASS: Component 1b -- Font size is 96pt ({font_size/12700:.0f}pt)")
                    else:
                        size_pt = font_size / 12700 if font_size else 'None'
                        print(f"FAIL: Component 1b -- Expected 96pt, found {size_pt}pt")

                    # Check bold
                    is_bold = runs[0].font.bold
                    if is_bold is True:
                        comp1_score += 0.06
                        print("PASS: Component 1c -- Font is bold")
                    else:
                        print(f"FAIL: Component 1c -- Expected bold=True, found {is_bold}")

                    # Check color #2C3E50
                    try:
                        if runs[0].font.color.type is not None:
                            color_rgb = str(runs[0].font.color.rgb).upper()
                            if color_rgb == '2C3E50':
                                comp1_score += 0.07
                                print(f"PASS: Component 1d -- Color is #2C3E50")
                            else:
                                print(f"FAIL: Component 1d -- Expected #2C3E50, found #{color_rgb}")
                        else:
                            print("FAIL: Component 1d -- No explicit color set")
                    except Exception as e:
                        print(f"FAIL: Component 1d -- Could not check color: {e}")

                # Check centered alignment
                from pptx.enum.text import PP_ALIGN
                if para.alignment == PP_ALIGN.CENTER:
                    comp1_score += 0.05
                    print("PASS: Component 1e -- Text is center-aligned")
                else:
                    print(f"FAIL: Component 1e -- Expected CENTER alignment, found {para.alignment}")

                # Check vertical position (y ~= 2 inches = 1828800 EMU)
                shape_top = shape_97.top
                if is_approx_equal(shape_top, Inches(2), tolerance=0.15):
                    comp1_score += 0.05
                    print(f"PASS: Component 1f -- Shape top at ~2in ({shape_top/914400:.2f}in)")
                else:
                    print(f"FAIL: Component 1f -- Expected top ~2in, found {shape_top/914400:.2f}in")

                total_score += comp1_score
                print(f"  Component 1 subtotal: {comp1_score:.2f}/{comp1_max}")

    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Thin underline rectangle at y=4in, 8in wide, #E74C3C (0.35 points)
    try:
        # Find thin rectangle-like shapes (height < 0.2 inches)
        underline_candidates = []
        for shape in shapes:
            if shape.height < Inches(0.2) and shape.width > Inches(4):
                underline_candidates.append(shape)

        if not underline_candidates:
            print("FAIL: Component 2 -- No thin wide shape (underline) found on slide 7")
        else:
            ul_shape = underline_candidates[0]
            comp2_score = 0.0
            comp2_max = 0.35

            # Check width (~8 inches = 7315200 EMU)
            if is_approx_equal(ul_shape.width, Inches(8), tolerance=0.1):
                comp2_score += 0.08
                print(f"PASS: Component 2a -- Width is ~8in ({ul_shape.width/914400:.2f}in)")
            else:
                print(f"FAIL: Component 2a -- Expected width ~8in, found {ul_shape.width/914400:.2f}in")

            # Check thinness (height should be very small, ~2pt or similar)
            if ul_shape.height < Inches(0.1):
                comp2_score += 0.05
                print(f"PASS: Component 2b -- Shape is thin ({ul_shape.height/914400:.3f}in)")
            else:
                print(f"FAIL: Component 2b -- Expected thin shape, found height {ul_shape.height/914400:.2f}in")

            # Check fill color #E74C3C
            try:
                fill = ul_shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    fill_color = str(fill.fore_color.rgb).upper()
                    if fill_color == 'E74C3C':
                        comp2_score += 0.10
                        print(f"PASS: Component 2c -- Fill color is #E74C3C")
                    else:
                        print(f"FAIL: Component 2c -- Expected fill #E74C3C, found #{fill_color}")
                else:
                    print(f"FAIL: Component 2c -- No solid fill on underline shape (type={fill.type})")
            except Exception as e:
                print(f"FAIL: Component 2c -- Could not check fill color: {e}")

            # Check vertical position (y ~= 4 inches = 3657600 EMU)
            if is_approx_equal(ul_shape.top, Inches(4), tolerance=0.1):
                comp2_score += 0.07
                print(f"PASS: Component 2d -- Top at ~4in ({ul_shape.top/914400:.2f}in)")
            else:
                print(f"FAIL: Component 2d -- Expected top ~4in, found {ul_shape.top/914400:.2f}in")

            # Check horizontal centering (left + width/2 should be near slide center)
            slide_center = prs.slide_width / 2
            shape_center = ul_shape.left + ul_shape.width / 2
            if is_approx_equal(shape_center, slide_center, tolerance=0.1):
                comp2_score += 0.05
                print(f"PASS: Component 2e -- Shape is horizontally centered")
            else:
                print(f"FAIL: Component 2e -- Shape not centered (center={shape_center/914400:.2f}in, slide center={slide_center/914400:.2f}in)")

            total_score += comp2_score
            print(f"  Component 2 subtotal: {comp2_score:.2f}/{comp2_max}")

    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: 'Client Satisfaction Rate' in 24pt #666666 centered at y=4.5in (0.30 points)
    try:
        matches_csr = find_text_shapes('client satisfaction rate')
        if not matches_csr:
            print("FAIL: Component 3 -- No shape with 'Client Satisfaction Rate' text found on slide 7")
        else:
            shape_csr = matches_csr[0]
            para = None
            for p in shape_csr.text_frame.paragraphs:
                if 'client satisfaction rate' in p.text.strip().lower():
                    para = p
                    break

            if para is None:
                print("FAIL: Component 3 -- Could not find paragraph with expected text")
            else:
                comp3_score = 0.0
                comp3_max = 0.30

                # Check text content
                runs = [r for r in para.runs if r.text.strip()]
                full_text = ''.join(r.text for r in runs).strip()
                if 'client satisfaction rate' in full_text.lower():
                    comp3_score += 0.05
                    print(f"PASS: Component 3a -- Text is '{full_text}'")
                else:
                    print(f"FAIL: Component 3a -- Expected 'Client Satisfaction Rate', found '{full_text}'")

                # Check font size (24pt = 304800 EMU)
                if runs:
                    font_size = runs[0].font.size
                    if font_size is not None and is_approx_equal(font_size, 304800):
                        comp3_score += 0.06
                        print(f"PASS: Component 3b -- Font size is 24pt ({font_size/12700:.0f}pt)")
                    else:
                        size_pt = font_size / 12700 if font_size else 'None'
                        print(f"FAIL: Component 3b -- Expected 24pt, found {size_pt}pt")

                    # Check color #666666
                    try:
                        if runs[0].font.color.type is not None:
                            color_rgb = str(runs[0].font.color.rgb).upper()
                            if color_rgb == '666666':
                                comp3_score += 0.07
                                print("PASS: Component 3c -- Color is #666666")
                            else:
                                print(f"FAIL: Component 3c -- Expected #666666, found #{color_rgb}")
                        else:
                            print("FAIL: Component 3c -- No explicit color set")
                    except Exception as e:
                        print(f"FAIL: Component 3c -- Could not check color: {e}")

                # Check centered alignment
                from pptx.enum.text import PP_ALIGN
                if para.alignment == PP_ALIGN.CENTER:
                    comp3_score += 0.05
                    print("PASS: Component 3d -- Text is center-aligned")
                else:
                    print(f"FAIL: Component 3d -- Expected CENTER alignment, found {para.alignment}")

                # Check vertical position (y ~= 4.5 inches = 4114800 EMU)
                if is_approx_equal(shape_csr.top, Inches(4.5), tolerance=0.1):
                    comp3_score += 0.07
                    print(f"PASS: Component 3e -- Shape top at ~4.5in ({shape_csr.top/914400:.2f}in)")
                else:
                    print(f"FAIL: Component 3e -- Expected top ~4.5in, found {shape_csr.top/914400:.2f}in")

                total_score += comp3_score
                print(f"  Component 3 subtotal: {comp3_score:.2f}/{comp3_max}")

    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
