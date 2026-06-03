"""
Initial Setup: Create a 10-slide Impact Report presentation with blank slide 7
Task ID: impress_design_085
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_085'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_title(slide, text):
    """Add a title-style text box at the top of a slide."""
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                 text, font_size=28, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50), alignment=PP_ALIGN.LEFT)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Use blank layout for all slides to avoid placeholder issues
    blank = prs.slide_layouts[6]  # use last layout (blank or near-blank)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(blank)
    add_text_box(slide1, Inches(1.5), Inches(2), Inches(7), Inches(1.5),
                 "2025 Annual Impact Report", font_size=40, bold=True,
                 color=RGBColor(0x2C, 0x3E, 0x50), alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(2), Inches(4), Inches(6), Inches(1),
                 "Greenfield Consulting Group", font_size=20,
                 color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(blank)
    add_title(slide2, "Executive Summary")
    add_text_box(slide2, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "This year marked a transformative period for Greenfield Consulting Group. "
                 "We expanded operations into three new markets, onboarded 42 enterprise clients, "
                 "and achieved record revenue growth of 31% year-over-year. Our commitment to "
                 "sustainability-driven strategy consulting continues to differentiate us in the market.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 3: Revenue Overview ---
    slide3 = prs.slides.add_slide(blank)
    add_title(slide3, "Revenue Overview")
    table_shape = slide3.shapes.add_table(5, 3, Inches(1.5), Inches(2), Inches(7), Inches(3))
    table = table_shape.table
    headers = ["Quarter", "Revenue ($M)", "Growth (%)"]
    data = [
        ["Q1 2025", "$12.4M", "28%"],
        ["Q2 2025", "$14.1M", "33%"],
        ["Q3 2025", "$15.8M", "35%"],
        ["Q4 2025", "$16.2M", "29%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 4: Market Expansion ---
    slide4 = prs.slides.add_slide(blank)
    add_title(slide4, "Market Expansion")
    add_text_box(slide4, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "New markets entered in 2025:\n"
                 "  - Southeast Asia (Singapore, Vietnam, Thailand)\n"
                 "  - Northern Europe (Denmark, Sweden, Finland)\n"
                 "  - South America (Brazil, Colombia)\n\n"
                 "Total addressable market increased by $2.3 billion.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 5: Team Growth ---
    slide5 = prs.slides.add_slide(blank)
    add_title(slide5, "Team Growth")
    add_text_box(slide5, Inches(0.8), Inches(1.8), Inches(8.4), Inches(2),
                 "Headcount grew from 186 to 274 employees across 8 offices worldwide. "
                 "We welcomed talent from McKinsey, BCG, Deloitte, and leading tech firms.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))
    tbl_shape = slide5.shapes.add_table(4, 2, Inches(2.5), Inches(4), Inches(5), Inches(2.5))
    tbl = tbl_shape.table
    dept_data = [
        ["Department", "New Hires"],
        ["Strategy", "34"],
        ["Technology", "28"],
        ["Operations", "26"],
    ]
    for r, row in enumerate(dept_data):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
            if r == 0:
                for run in tbl.cell(r, c).text_frame.paragraphs[0].runs:
                    run.font.bold = True

    # --- Slide 6: Sustainability Initiatives ---
    slide6 = prs.slides.add_slide(blank)
    add_title(slide6, "Sustainability Initiatives")
    add_text_box(slide6, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "Carbon neutral operations achieved in Q3 2025.\n"
                 "Partnered with 15 clients on ESG transformation projects.\n"
                 "Reduced business travel emissions by 42% through hybrid consulting model.\n"
                 "Published the Greenfield Sustainability Playbook, downloaded 12,000+ times.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 7: BLANK (agent must create statistics highlight here) ---
    slide7 = prs.slides.add_slide(blank)
    # Intentionally left completely blank

    # --- Slide 8: Client Portfolio ---
    slide8 = prs.slides.add_slide(blank)
    add_title(slide8, "Client Portfolio")
    add_text_box(slide8, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "Key enterprise clients added in 2025:\n"
                 "  - Meridian Healthcare Systems\n"
                 "  - Atlas Logistics International\n"
                 "  - NovaTech Solutions\n"
                 "  - Pinnacle Financial Group\n"
                 "  - Horizon Energy Partners\n\n"
                 "Client retention rate remains at an industry-leading level.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 9: Awards & Recognition ---
    slide9 = prs.slides.add_slide(blank)
    add_title(slide9, "Awards & Recognition")
    add_text_box(slide9, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "  - Forbes Top 50 Management Consulting Firms 2025\n"
                 "  - Financial Times Sustainability Leader Award\n"
                 "  - Glassdoor Best Places to Work (4.6/5.0 rating)\n"
                 "  - Harvard Business Review Case Study Feature",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 10: Looking Ahead ---
    slide10 = prs.slides.add_slide(blank)
    add_title(slide10, "Looking Ahead: 2026")
    add_text_box(slide10, Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5),
                 "Strategic priorities for 2026:\n"
                 "  - Launch AI-powered advisory practice\n"
                 "  - Expand to Middle East and Africa\n"
                 "  - Achieve B Corp certification\n"
                 "  - Grow revenue to $75M\n\n"
                 "We are positioned for another record-breaking year.",
                 font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
