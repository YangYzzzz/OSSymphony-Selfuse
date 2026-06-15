"""
Reward Script: Build a complete pricing page on slide 5 with three pricing cards
Task ID: impress_sales_052
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Three pricing card shapes with correct tier names (Basic, Professional, Enterprise)
  Component 2 (0.15): Correct feature counts per card (5, 8, 10)
  Component 3 (0.20): Card backgrounds (Basic=#F5F5F5, Pro=#FFFFFF w/ blue border, Enterprise=#F5F5F5)
  Component 4 (0.15): Three 'Get Started' buttons with blue fill (#2B6CB0) and white text
  Component 5 (0.10): 'MOST POPULAR' badge exists with blue fill and white text
  Component 6 (0.15): Pricing cards have shadows
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_052'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def check_shadows_via_xml(file_path, shape_names):
    """Check if specific shapes have outerShdw in their effectLst via XML."""
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
          'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    shadow_map = {}
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide5.xml') as f:
                root = ET.parse(f).getroot()
            for sp in root.findall('.//p:cSld//p:spTree//p:sp', ns):
                name_el = sp.find('.//p:nvSpPr/p:cNvPr', ns)
                if name_el is not None:
                    name = name_el.get('name', '')
                    shdw = sp.find('.//a:effectLst/a:outerShdw', ns)
                    shadow_map[name] = shdw is not None
    except Exception as e:
        print(f"ERROR: XML shadow check failed: {e}")
    return shadow_map


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # Slide 5 (0-indexed)

    # Classify shapes on slide 5 into cards, buttons, and badge
    cards = {}  # tier_name -> shape
    buttons = []
    badge = None

    for shape in slide.shapes:
        if not hasattr(shape, 'text') or not shape.text:
            continue
        text = shape.text.strip()

        # Check if this is a pricing card (contains tier name + price)
        if text.startswith('Basic') and '$29' in text:
            cards['Basic'] = shape
        elif text.startswith('Professional') and '$79' in text:
            cards['Professional'] = shape
        elif text.startswith('Enterprise') and '$199' in text:
            cards['Enterprise'] = shape
        elif text == 'Get Started':
            buttons.append(shape)
        elif 'MOST POPULAR' in text.upper():
            badge = shape

    # Component 1: Three pricing card shapes with correct tier names (0.25 points)
    try:
        found_cards = list(cards.keys())
        if len(found_cards) == 3 and all(t in cards for t in ['Basic', 'Professional', 'Enterprise']):
            print(f"PASS: Component 1 - All three pricing cards found: {found_cards} (0.25 pts)")
            total_score += 0.25
        elif len(found_cards) > 0:
            partial = 0.25 * len(found_cards) / 3.0
            print(f"PARTIAL: Component 1 - Found {len(found_cards)}/3 cards: {found_cards} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - No pricing cards found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Correct feature counts per card (0.15 points)
    # Basic: 5 features, Professional: 8 features, Enterprise: 10 features
    try:
        expected_features = {'Basic': 5, 'Professional': 8, 'Enterprise': 10}
        feature_checks_passed = 0
        for tier, expected_count in expected_features.items():
            if tier not in cards:
                print(f"FAIL: Component 2 - {tier} card missing, cannot check features")
                continue
            shape = cards[tier]
            # Features are lines after the empty paragraph (after 'per month')
            # Count non-empty paragraphs that look like feature items (indented or bulleted)
            if shape.has_text_frame:
                paras = shape.text_frame.paragraphs
                # Features start after the price/subtitle lines
                # Pattern: tier name, price, 'per month', empty, then features
                feature_lines = []
                found_empty = False
                for para in paras:
                    txt = para.text.strip()
                    if not txt:
                        found_empty = True
                        continue
                    if found_empty and txt:
                        feature_lines.append(txt)
                actual_count = len(feature_lines)
                if actual_count == expected_count:
                    print(f"PASS: Component 2 - {tier} has {actual_count} features (expected {expected_count})")
                    feature_checks_passed += 1
                else:
                    print(f"FAIL: Component 2 - {tier} has {actual_count} features, expected {expected_count}")
        if feature_checks_passed == 3:
            total_score += 0.15
            print(f"PASS: Component 2 - All feature counts correct (0.15 pts)")
        elif feature_checks_passed > 0:
            partial = 0.15 * feature_checks_passed / 3.0
            total_score += partial
            print(f"PARTIAL: Component 2 - {feature_checks_passed}/3 correct ({partial:.2f} pts)")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Card backgrounds and Professional border (0.20 points)
    try:
        bg_checks_passed = 0
        border_check_passed = False

        for tier, expected_bg in [('Basic', 'F5F5F5'), ('Professional', 'FFFFFF'), ('Enterprise', 'F5F5F5')]:
            if tier not in cards:
                print(f"FAIL: Component 3 - {tier} card missing")
                continue
            shape = cards[tier]
            try:
                fill = shape.fill
                if fill.type is not None:
                    actual_rgb = str(fill.fore_color.rgb)
                    if actual_rgb.upper() == expected_bg.upper():
                        print(f"PASS: Component 3 - {tier} background is #{actual_rgb}")
                        bg_checks_passed += 1
                    else:
                        print(f"FAIL: Component 3 - {tier} background is #{actual_rgb}, expected #{expected_bg}")
                else:
                    print(f"FAIL: Component 3 - {tier} has no solid fill")
            except Exception as e:
                print(f"ERROR: Component 3 - {tier} background check: {e}")

        # Check Professional card blue border (3pt = ~38100 EMU)
        if 'Professional' in cards:
            shape = cards['Professional']
            try:
                line = shape.line
                if line.fill.type is not None and line.fill.type == 1:  # SOLID
                    line_rgb = str(line.color.rgb).upper()
                    line_width = line.width
                    if line_rgb == '2B6CB0' and line_width is not None and line_width > 0:
                        border_check_passed = True
                        print(f"PASS: Component 3 - Professional border is blue #{line_rgb}, width={line_width}")
                    else:
                        print(f"FAIL: Component 3 - Professional border color={line_rgb}, width={line_width}")
                else:
                    print(f"FAIL: Component 3 - Professional card has no solid border")
            except Exception as e:
                print(f"ERROR: Component 3 - Professional border check: {e}")

        # Score: 0.15 for backgrounds (0.05 each) + 0.05 for border
        bg_score = 0.15 * bg_checks_passed / 3.0
        border_score = 0.05 if border_check_passed else 0.0
        comp3_score = bg_score + border_score
        total_score += comp3_score
        print(f"Component 3 total: {comp3_score:.2f}/0.20 pts (bg: {bg_checks_passed}/3, border: {'yes' if border_check_passed else 'no'})")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Three 'Get Started' buttons with blue fill and white text (0.15 points)
    try:
        valid_buttons = 0
        for btn in buttons:
            try:
                fill = btn.fill
                has_blue_fill = False
                has_white_text = False
                if fill.type is not None:
                    btn_rgb = str(fill.fore_color.rgb).upper()
                    if btn_rgb == '2B6CB0':
                        has_blue_fill = True

                if btn.has_text_frame:
                    for para in btn.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip() == 'Get Started':
                                try:
                                    if run.font.color.type is not None:
                                        txt_rgb = str(run.font.color.rgb).upper()
                                        if txt_rgb == 'FFFFFF':
                                            has_white_text = True
                                except:
                                    pass

                if has_blue_fill and has_white_text:
                    valid_buttons += 1
            except Exception:
                pass

        if valid_buttons == 3:
            print(f"PASS: Component 4 - 3 valid 'Get Started' buttons found (0.15 pts)")
            total_score += 0.15
        elif valid_buttons > 0:
            partial = 0.15 * valid_buttons / 3.0
            print(f"PARTIAL: Component 4 - {valid_buttons}/3 valid buttons ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - No valid 'Get Started' buttons found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: 'MOST POPULAR' badge exists (0.10 points)
    try:
        if badge is not None:
            badge_ok = False
            try:
                fill = badge.fill
                if fill.type is not None:
                    badge_rgb = str(fill.fore_color.rgb).upper()
                    badge_has_blue = badge_rgb == '2B6CB0'
                else:
                    badge_has_blue = False

                badge_has_white_text = False
                if badge.has_text_frame:
                    for para in badge.text_frame.paragraphs:
                        for run in para.runs:
                            if 'MOST POPULAR' in (run.text or '').upper():
                                try:
                                    if run.font.color.type is not None:
                                        txt_rgb = str(run.font.color.rgb).upper()
                                        if txt_rgb == 'FFFFFF':
                                            badge_has_white_text = True
                                except:
                                    pass

                if badge_has_blue and badge_has_white_text:
                    badge_ok = True
            except Exception:
                pass

            if badge_ok:
                print(f"PASS: Component 5 - 'MOST POPULAR' badge with blue fill & white text (0.10 pts)")
                total_score += 0.10
            else:
                print(f"PARTIAL: Component 5 - Badge found but styling incorrect (0.05 pts)")
                total_score += 0.05
        else:
            print(f"FAIL: Component 5 - No 'MOST POPULAR' badge found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    # Component 6: Pricing cards have shadows (0.15 points)
    try:
        shadow_map = check_shadows_via_xml(file_path, [])
        cards_with_shadow = 0
        for tier in ['Basic', 'Professional', 'Enterprise']:
            if tier not in cards:
                continue
            shape = cards[tier]
            shape_name = shape.name
            if shadow_map.get(shape_name, False):
                cards_with_shadow += 1
                print(f"PASS: Component 6 - {tier} card ({shape_name}) has shadow")
            else:
                print(f"FAIL: Component 6 - {tier} card ({shape_name}) has no shadow")

        if cards_with_shadow == 3:
            total_score += 0.15
            print(f"PASS: Component 6 - All 3 cards have shadows (0.15 pts)")
        elif cards_with_shadow > 0:
            partial = 0.15 * cards_with_shadow / 3.0
            total_score += partial
            print(f"PARTIAL: Component 6 - {cards_with_shadow}/3 cards have shadows ({partial:.2f} pts)")
        else:
            print(f"FAIL: Component 6 - No pricing cards have shadows")
    except Exception as e:
        print(f"ERROR: Component 6 - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
