"""
Initial Setup: Build a 10-slide sales presentation with slide 5 titled 'Pricing' but empty content area.
Task ID: impress_sales_052
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_052'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_titled_slide(prs, title_text):
    """Add a slide with title only, no body content."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2B, 0x6C, 0xB0)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "SalesForward 2025", "Annual Sales Strategy & Growth Plan")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Market Overview & Industry Trends",
        "Product Portfolio Update",
        "Target Customer Segments",
        "Pricing Strategy",
        "Go-to-Market Timeline",
        "Revenue Projections",
        "Team Structure & Responsibilities",
        "Q&A"
    ])

    # Slide 3: Market Overview
    add_content_slide(prs, "Market Overview", [
        "Total Addressable Market: $4.2B (growing 18% YoY)",
        "Key competitors: Acme Corp, NovaTech Solutions, CloudFirst Inc.",
        "Our current market share: 12% (up from 8% in 2024)",
        "Emerging opportunities in healthcare and fintech verticals",
        "Customer acquisition cost reduced by 22% through digital channels"
    ])

    # Slide 4: Product Portfolio
    add_content_slide(prs, "Product Portfolio", [
        "SalesForward Core - CRM & pipeline management",
        "SalesForward Analytics - Real-time dashboards & reporting",
        "SalesForward Connect - Omnichannel communication suite",
        "SalesForward AI - Predictive lead scoring & automation",
        "New: SalesForward Mobile - Field sales companion app"
    ])

    # Slide 5: Pricing (EMPTY content area - just title)
    add_blank_titled_slide(prs, "Pricing")

    # Slide 6: Customer Segments
    add_content_slide(prs, "Target Customer Segments", [
        "Small Business (1-50 employees): 40% of pipeline",
        "Mid-Market (51-500 employees): 35% of pipeline",
        "Enterprise (500+ employees): 25% of pipeline",
        "Focus verticals: Healthcare, Fintech, E-commerce, SaaS",
        "Average deal size increased to $24,500 in Q1 2025"
    ])

    # Slide 7: Go-to-Market Timeline
    add_content_slide(prs, "Go-to-Market Timeline", [
        "Q1 2025: Launch SalesForward Mobile beta",
        "Q2 2025: Enterprise pricing tier rollout",
        "Q3 2025: Partner channel program launch (50 partners)",
        "Q4 2025: International expansion - UK, Germany, Australia",
        "H1 2026: AI-powered automation features GA release"
    ])

    # Slide 8: Revenue Projections
    add_content_slide(prs, "Revenue Projections", [
        "2024 Actual: $18.7M ARR",
        "2025 Target: $28.5M ARR (+52% growth)",
        "Monthly recurring revenue target: $2.375M by Dec 2025",
        "Net revenue retention: 115%",
        "Gross margin target: 78%"
    ])

    # Slide 9: Team Structure
    add_content_slide(prs, "Sales Team Structure", [
        "VP of Sales: Jennifer Park",
        "Enterprise Sales (8 reps) - Director: Michael Torres",
        "Mid-Market Sales (12 reps) - Director: Rachel Kim",
        "SMB Sales (20 reps) - Director: David Okafor",
        "Sales Engineering (6 SEs) - Lead: Priya Sharma",
        "Sales Operations (4) - Manager: Carlos Rivera"
    ])

    # Slide 10: Thank You / Q&A
    add_title_slide(prs, "Thank You", "Questions & Discussion\ncontact@salesforward.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
