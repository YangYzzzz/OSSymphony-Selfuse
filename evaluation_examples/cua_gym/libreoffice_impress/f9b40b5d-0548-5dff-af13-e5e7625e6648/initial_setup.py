"""
Initial Setup: Create a 9-slide seminar presentation in LibreOffice Impress
Task ID: impress_el_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_el_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
    return slide


def add_section_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Advanced Data Visualization Techniques",
        "Faculty Seminar Series — Spring 2026\nDr. Elena Rodriguez, Department of Computer Science"
    )

    # Slide 2: Agenda
    add_content_slide(prs, "Seminar Agenda", [
        "Introduction to modern visualization paradigms",
        "Interactive dashboards with real-time data",
        "Case study: COVID-19 epidemiological tracking",
        "Hands-on exercise with Plotly and D3.js",
        "Panel discussion and Q&A",
    ])

    # Slide 3: Section header
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox = slide3.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Part 1: Foundations of Data Visualization"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Slide 4: Key Principles
    add_content_slide(prs, "Key Principles of Effective Visualization", [
        "Clarity: Reduce visual noise and emphasize patterns",
        "Accuracy: Faithful representation of underlying data",
        "Efficiency: Maximize data-ink ratio (Tufte, 2001)",
        "Accessibility: Color-blind safe palettes and alt text",
        "Context: Always provide axis labels, units, and sources",
    ])

    # Slide 5: Chart Types Overview
    add_content_slide(prs, "Choosing the Right Chart Type", [
        "Comparison: Bar charts, grouped bar, lollipop charts",
        "Distribution: Histograms, box plots, violin plots",
        "Composition: Stacked area, treemaps, sunburst diagrams",
        "Relationship: Scatter plots, bubble charts, heatmaps",
        "Temporal: Line charts, sparklines, Gantt charts",
    ])

    # Slide 6: Case Study intro
    add_content_slide(prs, "Case Study: Epidemiological Data Tracking", [
        "Data source: WHO Global Health Observatory (2020–2025)",
        "Challenge: Displaying 195 countries across 5 metrics",
        "Solution: Interactive choropleth with drill-down capability",
        "Tools used: Python (Pandas, Plotly), PostgreSQL, Dash",
        "Result: 40% faster anomaly detection by health officials",
    ])

    # Slide 7: Interactive Dashboards
    add_content_slide(prs, "Building Interactive Dashboards", [
        "Server-side rendering vs. client-side frameworks",
        "Real-time streaming with WebSocket connections",
        "Cross-filtering: Linking multiple views for exploration",
        "Performance: Aggregation strategies for 10M+ rows",
        "Deployment: Docker containers on university HPC cluster",
    ])

    # Slide 8: Hands-on Exercise
    add_content_slide(prs, "Hands-on Exercise: Plotly & D3.js", [
        "Exercise 1: Create a multi-series line chart with hover tooltips",
        "Exercise 2: Build a responsive scatter plot with zoom and pan",
        "Exercise 3: Design a geographic heatmap of research output",
        "All datasets available at: seminar.cs.example.edu/viz-data",
        "Submission deadline: Friday, April 10, 2026",
    ])

    # Slide 9: Summary & Q&A
    add_content_slide(prs, "Summary & Next Steps", [
        "Visualization is a critical tool for data-driven decisions",
        "Choose chart types based on the question, not the data",
        "Interactive tools empower non-technical stakeholders",
        "Next seminar: 'Machine Learning in Scientific Publishing'",
        "Contact: e.rodriguez@example.edu | Office: CS-412",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
