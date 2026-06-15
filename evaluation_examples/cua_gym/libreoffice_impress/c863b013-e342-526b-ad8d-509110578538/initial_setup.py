"""
Initial Setup: Quarterly Review presentation with TODO placeholders in slide notes
Task ID: impress_ndo_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 2025 Quarterly Business Review"
    slide1.placeholders[1].text = "Presented by: Strategy & Operations Division\nDecember 2025"
    slide1.notes_slide.notes_text_frame.text = (
        "Welcome everyone to the Q4 quarterly review.\n"
        "This presentation covers financial performance, product updates, and strategic initiatives."
    )

    # --- Slide 2: Financial Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Financial Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Revenue: $12.4M (+18% YoY)"
    body2.add_paragraph().text = "Gross Margin: 64.2%"
    body2.add_paragraph().text = "Operating Expenses: $7.8M"
    body2.add_paragraph().text = "Net Income: $1.9M"
    body2.add_paragraph().text = "Cash Position: $28.3M"
    slide2.notes_slide.notes_text_frame.text = (
        "Key financial highlights for Q4:\n"
        "- Revenue exceeded forecast by 3.2%\n"
        "- Margin improvement driven by automation savings\n"
        "[TODO: Add talking points]\n"
        "Remember to mention the FY2026 budget approval timeline."
    )

    # --- Slide 3: Customer Acquisition ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Customer Acquisition & Retention"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "New Enterprise Clients: 14"
    body3.add_paragraph().text = "SMB Sign-ups: 287"
    body3.add_paragraph().text = "Churn Rate: 2.1% (down from 3.4%)"
    body3.add_paragraph().text = "NPS Score: 72"
    body3.add_paragraph().text = "Expansion Revenue: $1.8M"
    slide3.notes_slide.notes_text_frame.text = (
        "Highlight the Meridian Healthcare and TechNova deals.\n"
        "Churn reduction attributed to the new onboarding program launched in September."
    )

    # --- Slide 4: Product Roadmap ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Product Roadmap Update"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "v3.2 Release: Shipped Oct 15"
    body4.add_paragraph().text = "API Gateway: Beta launched Nov 1"
    body4.add_paragraph().text = "Mobile App Redesign: 85% complete"
    body4.add_paragraph().text = "AI Assistant Feature: In development"
    body4.add_paragraph().text = "Security Audit: Passed SOC 2 Type II"
    slide4.notes_slide.notes_text_frame.text = (
        "Product milestones achieved this quarter:\n"
        "- v3.2 shipped on schedule with zero P1 bugs\n"
        "[TODO: Add talking points]\n"
        "Emphasize the AI assistant timeline and resource allocation for Q1."
    )

    # --- Slide 5: Team & Hiring ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Team Growth & Organizational Health"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Headcount: 142 (+23 this quarter)"
    body5.add_paragraph().text = "Open Positions: 18"
    body5.add_paragraph().text = "Employee Satisfaction: 4.3/5.0"
    body5.add_paragraph().text = "Voluntary Turnover: 8.2% annualized"
    body5.add_paragraph().text = "Training Hours per Employee: 24"
    slide5.notes_slide.notes_text_frame.text = (
        "Engineering team expanded significantly to support the AI initiative.\n"
        "HR flagged the need for additional recruiters in APAC region."
    )

    # --- Slide 6: Strategic Initiatives ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Strategic Initiatives for FY2026"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Initiative 1: International Expansion (EU & APAC)"
    body6.add_paragraph().text = "Initiative 2: Enterprise Platform Launch"
    body6.add_paragraph().text = "Initiative 3: Strategic Partnerships Program"
    body6.add_paragraph().text = "Initiative 4: Data Infrastructure Modernization"
    slide6.notes_slide.notes_text_frame.text = (
        "Strategic priorities for the upcoming fiscal year:\n"
        "- Board approved $4.5M for international expansion\n"
        "- Partnership discussions ongoing with Salesforce and AWS\n"
        "[TODO: Add talking points]\n"
        "Discuss timeline dependencies with the CFO before the board meeting."
    )

    # --- Slide 7: Risk Assessment ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Risk Assessment & Mitigation"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Market Risk: Competitor pricing pressure"
    body7.add_paragraph().text = "Regulatory: GDPR compliance for EU expansion"
    body7.add_paragraph().text = "Technical: Legacy system migration delays"
    body7.add_paragraph().text = "Talent: Senior engineering retention"
    slide7.notes_slide.notes_text_frame.text = (
        "Risk matrix updated with input from department heads.\n"
        "Compliance team recommends hiring a dedicated EU data protection officer by Q2."
    )

    # --- Slide 8: Thank You / Q&A ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[0])
    slide8.shapes.title.text = "Thank You"
    slide8.placeholders[1].text = "Questions & Discussion"
    slide8.notes_slide.notes_text_frame.text = (
        "Open the floor for questions.\n"
        "Key follow-up items will be tracked in the shared project tracker."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
