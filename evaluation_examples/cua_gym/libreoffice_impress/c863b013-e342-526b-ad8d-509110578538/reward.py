"""
Reward Script: Replace '[TODO: Add talking points]' placeholder in slide notes
Task ID: impress_ndo_029
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Slide 2 notes - placeholder replaced with correct text
  Component 2 (0.30): Slide 4 notes - placeholder replaced with correct text
  Component 3 (0.30): Slide 6 notes - placeholder replaced with correct text
  Component 4 (0.10): No residual placeholder text anywhere + other notes intact
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_029'

PLACEHOLDER = '[TODO: Add talking points]'
REPLACEMENT = 'Talking points finalized - see shared document.'

# Expected notes for slides that should NOT change (1-indexed keys)
UNCHANGED_NOTES = {
    1: 'Welcome everyone to the Q4 quarterly review.\nThis presentation covers financial performance, product updates, and strategic initiatives.',
    3: 'Highlight the Meridian Healthcare and TechNova deals.\nChurn reduction attributed to the new onboarding program launched in September.',
    5: 'Engineering team expanded significantly to support the AI initiative.\nHR flagged the need for additional recruiters in APAC region.',
    7: 'Risk matrix updated with input from department heads.\nCompliance team recommends hiring a dedicated EU data protection officer by Q2.',
    8: 'Open the floor for questions.\nKey follow-up items will be tracked in the shared project tracker.',
}

# Expected notes for changed slides (after replacement)
EXPECTED_NOTES = {
    2: 'Key financial highlights for Q4:\n- Revenue exceeded forecast by 3.2%\n- Margin improvement driven by automation savings\nTalking points finalized - see shared document.\nRemember to mention the FY2026 budget approval timeline.',
    4: 'Product milestones achieved this quarter:\n- v3.2 shipped on schedule with zero P1 bugs\nTalking points finalized - see shared document.\nEmphasize the AI assistant timeline and resource allocation for Q1.',
    6: 'Strategic priorities for the upcoming fiscal year:\n- Board approved $4.5M for international expansion\n- Partnership discussions ongoing with Salesforce and AWS\nTalking points finalized - see shared document.\nDiscuss timeline dependencies with the CFO before the board meeting.',
}


def persist_app_state():
    """Try to save any unsaved LibreOffice edits via Ctrl+S."""
    try:
        os.environ["DISPLAY"] = ":0"
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_notes(slide):
    """Get notes text from a slide, return empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)

    # Component 1: Slide 2 notes - placeholder replaced (0.30 points)
    # This checks that slide 2 notes contain the replacement text and NOT the placeholder
    try:
        notes_2 = get_notes(slides[1])  # 0-indexed
        has_replacement = REPLACEMENT in notes_2
        no_placeholder = PLACEHOLDER not in notes_2
        if has_replacement and no_placeholder:
            print(f"PASS: Component 1 -- Slide 2 notes contain replacement, no placeholder (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 -- Slide 2 notes: has_replacement={has_replacement}, no_placeholder={no_placeholder}")
            print(f"  Actual notes: {repr(notes_2[:200])}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Slide 4 notes - placeholder replaced (0.30 points)
    try:
        notes_4 = get_notes(slides[3])  # 0-indexed
        has_replacement = REPLACEMENT in notes_4
        no_placeholder = PLACEHOLDER not in notes_4
        if has_replacement and no_placeholder:
            print(f"PASS: Component 2 -- Slide 4 notes contain replacement, no placeholder (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 2 -- Slide 4 notes: has_replacement={has_replacement}, no_placeholder={no_placeholder}")
            print(f"  Actual notes: {repr(notes_4[:200])}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Slide 6 notes - placeholder replaced (0.30 points)
    try:
        notes_6 = get_notes(slides[5])  # 0-indexed
        has_replacement = REPLACEMENT in notes_6
        no_placeholder = PLACEHOLDER not in notes_6
        if has_replacement and no_placeholder:
            print(f"PASS: Component 3 -- Slide 6 notes contain replacement, no placeholder (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 3 -- Slide 6 notes: has_replacement={has_replacement}, no_placeholder={no_placeholder}")
            print(f"  Actual notes: {repr(notes_6[:200])}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: No residual placeholder anywhere + unchanged slides preserved (0.10 points)
    # This component verifies that:
    #   (a) no slide has the placeholder text remaining
    #   (b) unchanged slides (1,3,5,7,8) still have their original notes
    try:
        residual_found = False
        for i, slide in enumerate(slides):
            notes = get_notes(slide)
            if PLACEHOLDER in notes:
                print(f"FAIL: Component 4 -- Residual placeholder found on slide {i+1}")
                residual_found = True
                break

        unchanged_ok = True
        for slide_num, expected in UNCHANGED_NOTES.items():
            if slide_num - 1 < len(slides):
                actual = get_notes(slides[slide_num - 1])
                if actual != expected:
                    print(f"FAIL: Component 4 -- Slide {slide_num} notes changed unexpectedly")
                    print(f"  Expected: {repr(expected[:100])}")
                    print(f"  Actual:   {repr(actual[:100])}")
                    unchanged_ok = False
                    break

        if not residual_found and unchanged_ok:
            print(f"PASS: Component 4 -- No residual placeholders, unchanged slides intact (0.10 pts)")
            total_score += 0.10
        elif not residual_found:
            print(f"FAIL: Component 4 -- No residual placeholders but some unchanged slides were modified")
        # else: already printed fail for residual
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
