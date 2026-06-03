"""
Reward Script: Verify scatter plot on slide 6 with marketing spend vs revenue data
Task ID: impress_exec_042
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Scatter chart exists on slide 6
  - Component 2 (0.25): Chart title is 'Marketing ROI by Business Unit'
  - Component 3 (0.30): X values (marketing spend) match 6 expected data points
  - Component 4 (0.20): Y values (revenue) match 6 expected data points
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_042'

# Expected data points from the task instruction
# BU1: $500K spend, $4M rev; BU2: $800K, $6.5M; BU3: $1.2M, $9M
# BU4: $300K, $2.5M; BU5: $1.5M, $11M; BU6: $700K, $5.8M
EXPECTED_X = sorted([500000, 800000, 1200000, 300000, 1500000, 700000])
EXPECTED_Y_MAP = {
    500000: 4000000,
    800000: 6500000,
    1200000: 9000000,
    300000: 2500000,
    1500000: 11000000,
    700000: 5800000,
}


def persist_app_state(domain: str):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not installed")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # Slide 6 (0-indexed)

    # Find chart shape on slide 6
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Scatter chart exists on slide 6 (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            chart_type_val = chart.chart_type
            # XY_SCATTER is -4169
            is_scatter = (chart_type_val == -4169)
            if is_scatter:
                print(f"PASS: Component 1 — Scatter chart (XY_SCATTER) found on slide 6 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 — Chart found but type is {chart_type_val}, not XY_SCATTER (-4169)")
        else:
            print("FAIL: Component 1 — No chart found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if chart_shape is None:
        # No chart means no further checks possible
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Marketing ROI by Business Unit' (0.25 points)
    try:
        title_text = None
        try:
            title_text = chart.chart_title.text_frame.text.strip()
        except Exception:
            pass

        if title_text is not None and title_text == "Marketing ROI by Business Unit":
            print(f"PASS: Component 2 — Chart title matches exactly (0.25 pts)")
            total_score += 0.25
        elif title_text is not None and "marketing roi" in title_text.lower():
            print(f"PARTIAL: Component 2 — Title contains 'Marketing ROI' but not exact: '{title_text}' (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Expected title 'Marketing ROI by Business Unit', found: '{title_text}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Extract X and Y data from chart series
    actual_x = []
    actual_y = []
    try:
        if len(chart.series) > 0:
            series = chart.series[0]
            # Extract X values via XML
            ns = '{http://schemas.openxmlformats.org/drawingml/2006/chart}'
            xvals_elem = series._element.findall(f'.//{ns}xVal')
            for xv in xvals_elem:
                pts = xv.findall(f'.//{ns}v')
                actual_x = [float(p.text) for p in pts]

            # Extract Y values via XML
            yvals_elem = series._element.findall(f'.//{ns}yVal')
            for yv in yvals_elem:
                pts = yv.findall(f'.//{ns}v')
                actual_y = [float(p.text) for p in pts]

            print(f"  Extracted X values: {actual_x}")
            print(f"  Extracted Y values: {actual_y}")
        else:
            print("  No series found in chart")
    except Exception as e:
        print(f"  ERROR extracting chart data: {e}")

    # Component 3: X values (marketing spend) match expected (0.30 points)
    try:
        if len(actual_x) == 6:
            # Sort and compare with tolerance
            sorted_actual_x = sorted(actual_x)
            matches = 0
            for exp, act in zip(EXPECTED_X, sorted_actual_x):
                if abs(exp - act) / max(abs(exp), 1) <= 0.01:
                    matches += 1
            if matches == 6:
                print(f"PASS: Component 3 — All 6 X values (marketing spend) match (0.30 pts)")
                total_score += 0.30
            elif matches >= 4:
                partial = round(0.30 * matches / 6, 2)
                print(f"PARTIAL: Component 3 — {matches}/6 X values match ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 — Only {matches}/6 X values match. Expected (sorted): {EXPECTED_X}, Got: {sorted_actual_x}")
        elif len(actual_x) > 0:
            print(f"FAIL: Component 3 — Expected 6 X data points, found {len(actual_x)}")
        else:
            print("FAIL: Component 3 — No X values found in chart")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Y values (revenue) match expected (0.20 points)
    try:
        if len(actual_x) == 6 and len(actual_y) == 6:
            # Verify paired (x, y) values match the expected mapping
            matches = 0
            for ax, ay in zip(actual_x, actual_y):
                expected_y = EXPECTED_Y_MAP.get(ax)
                if expected_y is not None and abs(expected_y - ay) / max(abs(expected_y), 1) <= 0.01:
                    matches += 1
            if matches == 6:
                print(f"PASS: Component 4 — All 6 (x, y) data point pairs match (0.20 pts)")
                total_score += 0.20
            elif matches >= 4:
                partial = round(0.20 * matches / 6, 2)
                print(f"PARTIAL: Component 4 — {matches}/6 data pairs match ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — Only {matches}/6 data pairs match expected mapping")
        elif len(actual_y) > 0:
            print(f"FAIL: Component 4 — Expected 6 Y data points, found {len(actual_y)}")
        else:
            print("FAIL: Component 4 — No Y values found in chart")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
