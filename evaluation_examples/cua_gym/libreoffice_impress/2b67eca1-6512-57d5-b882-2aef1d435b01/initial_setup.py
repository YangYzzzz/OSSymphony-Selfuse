"""
Initial Setup: Marketing Review presentation with 8 slides. Slide 6 has title only.
Task ID: impress_exec_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_textbox(slide, prs, text, top=Inches(0.4), font_size=Pt(28), bold=True, color=RGBColor(0x1B, 0x3A, 0x5C)):
    """Add a styled title textbox to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_body_textbox(slide, text_lines, top=Inches(1.5), left=Inches(0.8), width=Inches(8.4), height=Inches(4.5)):
    """Add body text with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---------- Slide 1: Title Slide ----------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Q1 2025 Marketing Review"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Prepared by the Strategic Marketing Division  |  April 2025"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # ---------- Slide 2: Executive Summary ----------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide2, prs, "Executive Summary")
    add_body_textbox(slide2, [
        "Total marketing spend in Q1 2025 reached $5.0M across all business units,",
        "representing a 12% increase over Q4 2024. Revenue attribution models indicate",
        "a blended ROI of 7.2x, with significant variance across units.",
        "",
        "Key highlights:",
        "  - Digital channels drove 68% of qualified leads",
        "  - Brand awareness scores improved by 15 points (NPS: 42 to 57)",
        "  - Customer acquisition cost decreased 8% quarter-over-quarter",
        "  - BU5 achieved highest ROI at 7.3x marketing spend",
    ])

    # ---------- Slide 3: Campaign Performance ----------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide3, prs, "Campaign Performance Overview")

    # Add a table with campaign data
    rows, cols = 7, 5
    tbl_shape = slide3.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.0))
    table = tbl_shape.table
    headers = ["Campaign", "Channel", "Spend ($K)", "Leads Generated", "Conv. Rate"]
    data = [
        ["Spring Product Launch", "Paid Search", "320", "4,250", "3.8%"],
        ["Brand Refresh 2025", "Social Media", "185", "2,800", "2.1%"],
        ["Enterprise Outreach", "Email", "95", "1,620", "5.2%"],
        ["Partner Co-Marketing", "Events", "210", "980", "4.6%"],
        ["Content Series: AI Trends", "Organic", "45", "3,100", "1.9%"],
        ["Retargeting Wave 3", "Display", "150", "2,340", "3.3%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # ---------- Slide 4: Digital Marketing Metrics ----------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide4, prs, "Digital Marketing Metrics")
    add_body_textbox(slide4, [
        "Website Traffic: 1.2M sessions (+18% QoQ)",
        "Organic Search: 480K sessions, 40% of total traffic",
        "Paid Search CTR: 4.2% (industry avg: 3.1%)",
        "Social Media Engagement Rate: 6.8% across platforms",
        "Email Open Rate: 28.4% | Click-Through Rate: 4.1%",
        "Landing Page Conversion: 3.6% (up from 2.9% in Q4)",
        "",
        "Top performing content: 'AI in Marketing 2025' whitepaper",
        "  - 12,400 downloads in 6 weeks",
        "  - Generated 890 MQLs directly attributed",
    ])

    # ---------- Slide 5: Budget Overview ----------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide5, prs, "Q1 Budget Allocation & Actuals")

    rows2, cols2 = 8, 4
    tbl_shape2 = slide5.shapes.add_table(rows2, cols2, Inches(1), Inches(1.5), Inches(8), Inches(4.0))
    table2 = tbl_shape2.table
    headers2 = ["Business Unit", "Budget ($K)", "Actual ($K)", "Variance"]
    budget_data = [
        ["BU1 - Consumer Products", "520", "500", "-3.8%"],
        ["BU2 - Enterprise Solutions", "810", "800", "-1.2%"],
        ["BU3 - Cloud Services", "1,180", "1,200", "+1.7%"],
        ["BU4 - Emerging Markets", "320", "300", "-6.3%"],
        ["BU5 - Digital Platforms", "1,450", "1,500", "+3.4%"],
        ["BU6 - Professional Services", "720", "700", "-2.8%"],
        ["TOTAL", "5,000", "5,000", "0.0%"],
    ]
    for c, h in enumerate(headers2):
        cell = table2.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1B, 0x3A, 0x5C)
    for r, row_data in enumerate(budget_data, 1):
        for c, val in enumerate(row_data):
            table2.cell(r, c).text = val
        if r == 7:  # Total row bold
            for c in range(cols2):
                for run in table2.cell(r, c).text_frame.paragraphs[0].runs:
                    run.font.bold = True

    # ---------- Slide 6: Marketing Efficiency (TITLE ONLY, no content) ----------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide6, prs, "Marketing Efficiency")
    # NO chart, NO scatter plot, NO data content - this is the task target

    # ---------- Slide 7: Competitive Analysis ----------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide7, prs, "Competitive Landscape")
    add_body_textbox(slide7, [
        "Market Position: #2 in share of voice (up from #3 in Q4 2024)",
        "",
        "Competitor A: Increased paid media spend by ~25%, focusing on video",
        "Competitor B: Launched aggressive referral program in Feb 2025",
        "Competitor C: Expanded into APAC with localized content strategy",
        "",
        "Our differentiation: Thought leadership content and partner ecosystem",
        "continue to be our strongest competitive advantages, particularly",
        "in the enterprise segment where trust signals matter most.",
    ])

    # ---------- Slide 8: Next Steps ----------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_textbox(slide8, prs, "Q2 2025 Priorities & Next Steps")
    add_body_textbox(slide8, [
        "1. Scale BU5 playbook to BU2 and BU3 (projected +15% ROI lift)",
        "2. Launch ABM pilot targeting top 50 enterprise accounts",
        "3. Increase video content budget by 30% based on Q1 engagement data",
        "4. Deploy marketing attribution model v2.0 (multi-touch)",
        "5. Hire 2 additional data analysts for marketing operations",
        "",
        "Timeline: Initiative kickoffs scheduled for April 14-18, 2025",
        "Budget request: Additional $450K for Q2 expansion initiatives",
        "Executive review: May 5, 2025 with full Q2 mid-quarter update",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
