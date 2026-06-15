"""
Reward Script: Apply sunshine yellow background to slides 2, 5, 8 and update slide 1 title
Task ID: osworld_impress_conditional_bg_image_012
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 background == #FFF44F (0.25 pts)
  Component 2: Slide 5 background == #FFF44F (0.25 pts)
  Component 3: Slide 8 background == #FFF44F (0.25 pts)
  Component 4: Slide 1 title text == 'Visual Team Report — Spring Edition' (0.25 pts)
  Total: 1.0
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_012'

TARGET_BG_COLOR = 'FFF44F'  # sunshine yellow #FFF44F
TARGET_TITLE = 'Visual Team Report \u2014 Spring Edition'  # em-dash

# Slides that should have yellow background (1-indexed in task, 0-indexed in API)
YELLOW_SLIDES = [1, 4, 7]  # slide indices 0-based for slides 2, 5, 8


def get_slide_background_rgb(slide):
    """Return the background RGB hex string for a slide, or None."""
    fill = slide.background.fill
    if fill.type == 1:  # solid fill directly on slide
        return str(fill.fore_color.rgb)
    elif fill.type == 5:  # inherited from master/layout
        # Check layout first
        layout_fill = slide.slide_layout.background.fill
        if layout_fill.type == 1:
            return str(layout_fill.fore_color.rgb)
        # Fall back to slide master
        master_fill = slide.slide_layout.slide_master.background.fill
        if master_fill.type == 1:
            return str(master_fill.fore_color.rgb)
        return None
    return None


def get_slide1_title_text(prs):
    """Return the title text of slide 1 (first non-empty text in first shape)."""
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 9 slides
    if len(prs.slides) != 9:
        print(f"CRITICAL: Expected 9 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 (index 1) background is #FFF44F (0.25 points)
    try:
        bg = get_slide_background_rgb(prs.slides[1])
        if bg and bg.upper() == TARGET_BG_COLOR:
            print(f"PASS: Component 1 — Slide 2 background is #{TARGET_BG_COLOR} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Slide 2 background expected #{TARGET_BG_COLOR}, found {bg!r}")
    except Exception as e:
        print(f"ERROR: Component 1 (Slide 2 background) — {e}")

    # Component 2: Slide 5 (index 4) background is #FFF44F (0.25 points)
    try:
        bg = get_slide_background_rgb(prs.slides[4])
        if bg and bg.upper() == TARGET_BG_COLOR:
            print(f"PASS: Component 2 — Slide 5 background is #{TARGET_BG_COLOR} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Slide 5 background expected #{TARGET_BG_COLOR}, found {bg!r}")
    except Exception as e:
        print(f"ERROR: Component 2 (Slide 5 background) — {e}")

    # Component 3: Slide 8 (index 7) background is #FFF44F (0.25 points)
    try:
        bg = get_slide_background_rgb(prs.slides[7])
        if bg and bg.upper() == TARGET_BG_COLOR:
            print(f"PASS: Component 3 — Slide 8 background is #{TARGET_BG_COLOR} (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Slide 8 background expected #{TARGET_BG_COLOR}, found {bg!r}")
    except Exception as e:
        print(f"ERROR: Component 3 (Slide 8 background) — {e}")

    # Component 4: Slide 1 title text updated to 'Visual Team Report — Spring Edition' (0.25 points)
    try:
        title_text = get_slide1_title_text(prs)
        if title_text and title_text.strip() == TARGET_TITLE:
            print(f"PASS: Component 4 — Slide 1 title is '{TARGET_TITLE}' (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Slide 1 title expected '{TARGET_TITLE}', found {title_text!r}")
    except Exception as e:
        print(f"ERROR: Component 4 (Slide 1 title) — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
