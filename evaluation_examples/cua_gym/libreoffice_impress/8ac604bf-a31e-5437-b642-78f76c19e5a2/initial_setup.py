"""
Initial Setup: 9-slide team report presentation with white backgrounds on all slides.
Task ID: osworld_impress_conditional_bg_image_012
Domain: libreoffice_impress

The agent must:
  1. Apply #FFF44F background to slides 2, 5, and 8 (person photo slides)
  2. Update slide 1 title from 'Team Report' to 'Visual Team Report — Spring Edition'
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw
import io

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_conditional_bg_image_012'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def make_person_photo_bytes(width=400, height=300, bg_color=(200, 210, 220),
                             skin_color=(220, 180, 140), label="Team Member"):
    """Generate a simple placeholder person-photo image as bytes."""
    img = Image.new('RGB', (width, height), color=bg_color)
    draw = ImageDraw.Draw(img)
    # Draw a simple figure: circle head + rectangle body
    cx, cy = width // 2, height // 2
    # Head
    head_r = width // 8
    draw.ellipse([cx - head_r, cy - head_r - 40, cx + head_r, cy + head_r - 40],
                 fill=skin_color, outline=(160, 120, 80), width=2)
    # Body
    body_w, body_h = width // 4, height // 3
    draw.rectangle([cx - body_w // 2, cy - 10, cx + body_w // 2, cy + body_h],
                   fill=(100, 130, 180), outline=(70, 90, 140), width=2)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def make_chart_placeholder_bytes(width=400, height=250, label="Q1 Performance"):
    """Generate a simple bar-chart placeholder image as bytes."""
    img = Image.new('RGB', (width, height), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)
    bars = [
        ((60, 180, 100, 230), (70, 130, 90)),
        ((120, 140, 160, 230), (70, 130, 90)),
        ((180, 100, 220, 230), (70, 130, 90)),
        ((240, 160, 280, 230), (70, 130, 90)),
        ((300, 80, 340, 230), (70, 130, 90)),
    ]
    colors = [(70, 130, 180), (100, 180, 100), (200, 100, 80), (180, 140, 60), (140, 80, 180)]
    for i, (rect, _) in enumerate(bars):
        draw.rectangle(rect, fill=colors[i % len(colors)])
    draw.line([(50, 230), (370, 230)], fill=(80, 80, 80), width=2)
    draw.line([(50, 40), (50, 230)], fill=(80, 80, 80), width=2)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def set_white_background(slide):
    """Set a solid white background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_text(slide, title_text, font_size=28, bold=True,
                   font_color=RGBColor(0x1F, 0x38, 0x64)):
    """Add a title text box to slide."""
    # Use fixed slide width (10 inches = 9144000 EMU)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold = bold
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color


def add_body_text(slide, body_text, top_offset=Inches(1.5), font_size=14,
                  color=RGBColor(0x33, 0x33, 0x33)):
    """Add a body text box to slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), top_offset, Inches(9.0), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body_text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # ── Slide 1: Title slide ──────────────────────────────────────────────────
    sl1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_white_background(sl1)
    # Title text box — initial state: "Team Report"
    txTitle = sl1.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(8.0), Inches(1.5))
    tf = txTitle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Team Report"
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # Subtitle
    txSub = sl1.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.8))
    tf2 = txSub.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Quarterly Review | Engineering Division"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

    # ── Slide 2: Person photo — white background (NOT yellow) ────────────────
    sl2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl2)
    add_title_text(sl2, "Meet the Team — Sarah Chen, Engineering Lead")
    photo_buf2 = make_person_photo_bytes(380, 280, bg_color=(210, 220, 230),
                                         skin_color=(240, 195, 155), label="Sarah Chen")
    sl2.shapes.add_picture(photo_buf2, Inches(0.8), Inches(1.6), Inches(4.2), Inches(3.1))
    # Bio text box
    bio2 = sl2.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.0), Inches(3.0))
    tf_bio2 = bio2.text_frame
    tf_bio2.word_wrap = True
    lines2 = [
        ("Name:", "Sarah Chen"),
        ("Role:", "Engineering Lead"),
        ("Department:", "Platform Engineering"),
        ("Years at Company:", "7"),
        ("Projects:", "Cloud Migration, CI/CD Overhaul"),
        ("Quote:", '"Building reliable systems is my passion."'),
    ]
    first2 = True
    for label, val in lines2:
        p_new = tf_bio2.paragraphs[0] if first2 else tf_bio2.add_paragraph()
        first2 = False
        r_lbl = p_new.add_run()
        r_lbl.text = f"{label} "
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(12)
        r_lbl.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        r_val = p_new.add_run()
        r_val.text = val
        r_val.font.size = Pt(12)
        r_val.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 3: Metrics / Data ───────────────────────────────────────────────
    sl3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl3)
    add_title_text(sl3, "Q1 Key Metrics — Engineering Division")
    chart_buf3 = make_chart_placeholder_bytes(500, 300, "Q1 Performance")
    sl3.shapes.add_picture(chart_buf3, Inches(0.5), Inches(1.5), Inches(5.5), Inches(3.5))
    metrics_box = sl3.shapes.add_textbox(Inches(6.3), Inches(1.5), Inches(3.2), Inches(3.5))
    tf_m = metrics_box.text_frame
    tf_m.word_wrap = True
    metric_items = [
        "Deployments: 148",
        "Uptime: 99.97%",
        "P1 Incidents: 2",
        "Features Shipped: 34",
        "Avg. Build Time: 4.2 min",
    ]
    for i, item in enumerate(metric_items):
        para = tf_m.paragraphs[0] if i == 0 else tf_m.add_paragraph()
        r = para.add_run()
        r.text = f"• {item}"
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    # ── Slide 4: Project Summary ──────────────────────────────────────────────
    sl4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl4)
    add_title_text(sl4, "Project Highlights — Q1 2025")
    projects = [
        ("Cloud Migration Phase 2", "Migrated 12 services to AWS; 30% cost reduction."),
        ("CI/CD Pipeline Overhaul", "Reduced deployment time from 45 min to 8 min."),
        ("Data Platform Upgrade", "Upgraded Kafka to v3.6; improved throughput by 40%."),
        ("Security Audit Response", "Resolved 18 critical vulnerabilities; zero P0 issues."),
    ]
    top = Inches(1.5)
    for proj_title, proj_desc in projects:
        box = sl4.shapes.add_textbox(Inches(0.5), top, Inches(9.0), Inches(0.75))
        tf_p = box.text_frame
        tf_p.word_wrap = True
        p_head = tf_p.paragraphs[0]
        r_head = p_head.add_run()
        r_head.text = proj_title
        r_head.font.bold = True
        r_head.font.size = Pt(14)
        r_head.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        p_desc = tf_p.add_paragraph()
        r_desc = p_desc.add_run()
        r_desc.text = proj_desc
        r_desc.font.size = Pt(12)
        r_desc.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        top = top + Inches(1.05)

    # ── Slide 5: Person photo — white background (NOT yellow) ────────────────
    sl5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl5)
    add_title_text(sl5, "Meet the Team — Marcus Johnson, DevOps Manager")
    photo_buf5 = make_person_photo_bytes(380, 280, bg_color=(220, 225, 215),
                                         skin_color=(180, 140, 100), label="Marcus Johnson")
    sl5.shapes.add_picture(photo_buf5, Inches(0.8), Inches(1.6), Inches(4.2), Inches(3.1))
    bio5 = sl5.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.0), Inches(3.0))
    tf_bio5 = bio5.text_frame
    tf_bio5.word_wrap = True
    lines5 = [
        ("Name:", "Marcus Johnson"),
        ("Role:", "DevOps Manager"),
        ("Department:", "Infrastructure"),
        ("Years at Company:", "5"),
        ("Projects:", "Kubernetes Migration, SRE Practice"),
        ("Quote:", '"Automate everything that can be automated."'),
    ]
    first5 = True
    for label, val in lines5:
        p_new = tf_bio5.paragraphs[0] if first5 else tf_bio5.add_paragraph()
        first5 = False
        r_lbl = p_new.add_run()
        r_lbl.text = f"{label} "
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(12)
        r_lbl.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        r_val = p_new.add_run()
        r_val.text = val
        r_val.font.size = Pt(12)
        r_val.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 6: Team Goals ───────────────────────────────────────────────────
    sl6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl6)
    add_title_text(sl6, "Q2 2025 Goals & Priorities")
    goals = [
        "Launch self-service developer portal by April 15",
        "Achieve 99.99% uptime SLA for Tier-1 services",
        "Complete SOC 2 Type II audit preparation",
        "Onboard 4 new engineers — backend & SRE tracks",
        "Reduce mean time to recovery (MTTR) by 25%",
        "Implement FinOps tooling for cost visibility",
    ]
    top_g = Inches(1.6)
    for goal in goals:
        g_box = sl6.shapes.add_textbox(Inches(0.7), top_g, Inches(8.6), Inches(0.55))
        tf_g = g_box.text_frame
        p_g = tf_g.paragraphs[0]
        r_g = p_g.add_run()
        r_g.text = f"➤  {goal}"
        r_g.font.size = Pt(14)
        r_g.font.color.rgb = RGBColor(0x20, 0x60, 0x30)
        top_g = top_g + Inches(0.7)

    # ── Slide 7: Data Table ───────────────────────────────────────────────────
    sl7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl7)
    add_title_text(sl7, "Team Capacity & Allocation")
    table_shape = sl7.shapes.add_table(6, 4, Inches(0.5), Inches(1.5),
                                        Inches(9.0), Inches(3.0))
    tbl = table_shape.table
    headers_tbl = ["Team Member", "Role", "Allocation", "Status"]
    rows_tbl = [
        ["Sarah Chen", "Engineering Lead", "100%", "Active"],
        ["Marcus Johnson", "DevOps Manager", "100%", "Active"],
        ["Priya Patel", "Backend Engineer", "80%", "On-loan"],
        ["David Kim", "SRE", "100%", "Active"],
        ["Lena Muller", "Data Engineer", "60%", "Part-time"],
    ]
    for c_idx, h in enumerate(headers_tbl):
        cell = tbl.cell(0, c_idx)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1F, 0x38, 0x64)
    for r_idx, row_data in enumerate(rows_tbl, 1):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # ── Slide 8: Person photo — white background (NOT yellow) ────────────────
    sl8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl8)
    add_title_text(sl8, "Meet the Team — Priya Patel, Backend Engineer")
    photo_buf8 = make_person_photo_bytes(380, 280, bg_color=(225, 215, 230),
                                         skin_color=(210, 170, 130), label="Priya Patel")
    sl8.shapes.add_picture(photo_buf8, Inches(0.8), Inches(1.6), Inches(4.2), Inches(3.1))
    bio8 = sl8.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(4.0), Inches(3.0))
    tf_bio8 = bio8.text_frame
    tf_bio8.word_wrap = True
    lines8 = [
        ("Name:", "Priya Patel"),
        ("Role:", "Backend Engineer"),
        ("Department:", "Platform Engineering"),
        ("Years at Company:", "3"),
        ("Projects:", "API Gateway Redesign, Auth Service"),
        ("Quote:", '"Great code reads like a good story."'),
    ]
    first8 = True
    for label, val in lines8:
        p_new = tf_bio8.paragraphs[0] if first8 else tf_bio8.add_paragraph()
        first8 = False
        r_lbl = p_new.add_run()
        r_lbl.text = f"{label} "
        r_lbl.font.bold = True
        r_lbl.font.size = Pt(12)
        r_lbl.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        r_val = p_new.add_run()
        r_val.text = val
        r_val.font.size = Pt(12)
        r_val.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 9: Closing / Summary ────────────────────────────────────────────
    sl9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_background(sl9)
    tx_close = sl9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(1.5))
    tf_close = tx_close.text_frame
    p_close = tf_close.paragraphs[0]
    p_close.alignment = PP_ALIGN.CENTER
    r_close = p_close.add_run()
    r_close.text = "Thank You"
    r_close.font.bold = True
    r_close.font.size = Pt(42)
    r_close.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    tx_sub9 = sl9.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(1.0))
    tf_sub9 = tx_sub9.text_frame
    p_sub9 = tf_sub9.paragraphs[0]
    p_sub9.alignment = PP_ALIGN.CENTER
    r_sub9 = p_sub9.add_run()
    r_sub9.text = "Engineering Division | Q1 2025 Report"
    r_sub9.font.size = Pt(18)
    r_sub9.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Slides: {len(prs.slides)} total')
    print('  Slide 1: Title — "Team Report" (white bg)')
    print('  Slide 2: Sarah Chen photo (white bg — task target)')
    print('  Slide 3: Q1 Metrics chart (white bg)')
    print('  Slide 4: Project Highlights (white bg)')
    print('  Slide 5: Marcus Johnson photo (white bg — task target)')
    print('  Slide 6: Q2 Goals (white bg)')
    print('  Slide 7: Team Capacity table (white bg)')
    print('  Slide 8: Priya Patel photo (white bg — task target)')
    print('  Slide 9: Closing (white bg)')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
