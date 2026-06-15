"""
Initial Setup: Create a Geometry presentation with slide 2 'Basic Shapes' having no shapes yet.
Task ID: impress_ndo_061
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_061'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Geometry Lessons"
    slide1.placeholders[1].text = "A Visual Guide to Shapes and Figures"

    # --- Slide 2: Basic Shapes (Title Only - NO shapes yet) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add a title text box at the top
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Basic Shapes"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    # NO shapes on this slide - that's the task

    # --- Slide 3: Properties of Triangles ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf3t = txBox3_title.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Properties of Triangles"
    p3t.alignment = PP_ALIGN.CENTER
    run3t = p3t.runs[0]
    run3t.font.size = Pt(32)
    run3t.font.bold = True
    run3t.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txBox3 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "An isosceles triangle has two sides of equal length."
    run3 = p3.runs[0]
    run3.font.size = Pt(20)
    run3.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    p3b = tf3.add_paragraph()
    p3b.text = "The angles opposite the equal sides are also equal."
    run3b = p3b.runs[0]
    run3b.font.size = Pt(20)
    run3b.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    p3c = tf3.add_paragraph()
    p3c.text = "The sum of interior angles in any triangle is always 180 degrees."
    run3c = p3c.runs[0]
    run3c.font.size = Pt(20)
    run3c.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    # --- Slide 4: Quadrilaterals Overview ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf4t = txBox4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Quadrilaterals Overview"
    p4t.alignment = PP_ALIGN.CENTER
    run4t = p4t.runs[0]
    run4t.font.size = Pt(32)
    run4t.font.bold = True
    run4t.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    txBox4 = slide4.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    items = [
        "Square: All sides equal, all angles 90 degrees",
        "Rectangle: Opposite sides equal, all angles 90 degrees",
        "Rhombus: All sides equal, opposite angles equal",
        "Trapezoid: One pair of parallel sides",
    ]
    for i, item in enumerate(items):
        if i == 0:
            p4 = tf4.paragraphs[0]
        else:
            p4 = tf4.add_paragraph()
        p4.text = item
        run4 = p4.runs[0]
        run4.font.size = Pt(18)
        run4.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
