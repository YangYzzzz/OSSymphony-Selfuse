"""
Reward Script: Insert isosceles triangle on slide 2 with specific properties
Task ID: impress_ndo_061
Domain: libreoffice_impress
Scoring:
  Component 1 (0.20): Isosceles triangle shape exists on slide 2
  Component 2 (0.20): Size ~6cm wide x 8cm tall
  Component 3 (0.15): Centered on the slide
  Component 4 (0.15): Fill color #1ABC9C
  Component 5 (0.15): Border 1.5pt solid #16A085
  Component 6 (0.15): Text 'A' in 24pt bold white
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_061'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.dml.color import RGBColor
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # Slide 2 (0-indexed)

    # Find isosceles triangle shapes on slide 2
    triangle = None
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.ISOSCELES_TRIANGLE:
                    triangle = shape
                    break
            except Exception:
                pass
            # Fallback: check shape name
            if triangle is None and 'triangle' in shape.name.lower():
                triangle = shape
                break

    # Component 1: Isosceles triangle shape exists on slide 2 (0.20 points)
    try:
        if triangle is not None:
            print(f"PASS: Component 1 — Isosceles triangle found: '{triangle.name}' (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — No isosceles triangle shape found on slide 2")
            # Without the triangle, nothing else can be checked
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Size approximately 6cm wide x 8cm tall (0.20 points)
    # 1 cm = 360000 EMU; 6cm = 2160000, 8cm = 2880000
    try:
        expected_w = 2160000  # 6cm in EMU
        expected_h = 2880000  # 8cm in EMU
        actual_w = triangle.width
        actual_h = triangle.height
        # Use 10% tolerance for size
        w_ok = abs(actual_w - expected_w) / expected_w <= 0.10
        h_ok = abs(actual_h - expected_h) / expected_h <= 0.10
        if w_ok and h_ok:
            print(f"PASS: Component 2 — Size {actual_w/360000:.2f}cm x {actual_h/360000:.2f}cm (0.20 pts)")
            total_score += 0.20
        elif w_ok or h_ok:
            print(f"PARTIAL: Component 2 — Width {'OK' if w_ok else 'WRONG'} ({actual_w/360000:.2f}cm), "
                  f"Height {'OK' if h_ok else 'WRONG'} ({actual_h/360000:.2f}cm) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Size {actual_w/360000:.2f}cm x {actual_h/360000:.2f}cm, "
                  f"expected ~6cm x ~8cm")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Centered on the slide (0.15 points)
    try:
        slide_w = prs.slide_width
        slide_h = prs.slide_height
        shape_cx = triangle.left + triangle.width / 2
        shape_cy = triangle.top + triangle.height / 2
        slide_cx = slide_w / 2
        slide_cy = slide_h / 2
        # Use 5% of slide dimension as tolerance
        x_tol = slide_w * 0.05
        y_tol = slide_h * 0.05
        x_ok = abs(shape_cx - slide_cx) <= x_tol
        y_ok = abs(shape_cy - slide_cy) <= y_tol
        if x_ok and y_ok:
            print(f"PASS: Component 3 — Shape centered on slide (0.15 pts)")
            total_score += 0.15
        elif x_ok:
            print(f"PARTIAL: Component 3 — Horizontally centered, vertical offset "
                  f"{abs(shape_cy - slide_cy)/360000:.2f}cm (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 3 — Shape not centered. "
                  f"Shape center=({shape_cx}, {shape_cy}), Slide center=({slide_cx}, {slide_cy})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Fill color #1ABC9C (0.15 points)
    try:
        fill = triangle.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            fill_rgb = str(fill.fore_color.rgb).upper()
            if fill_rgb == '1ABC9C':
                print(f"PASS: Component 4 — Fill color #{fill_rgb} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Fill color #{fill_rgb}, expected #1ABC9C")
        else:
            print(f"FAIL: Component 4 — Fill type is {fill.type}, expected solid fill")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Border 1.5pt solid #16A085 (0.15 points)
    try:
        line = triangle.line
        line_width = line.width  # in EMU; 1pt = 12700 EMU; 1.5pt = 19050
        line_color = None
        try:
            line_color = str(line.color.rgb).upper()
        except Exception:
            pass

        width_ok = line_width is not None and abs(line_width - 19050) / 19050 <= 0.10
        color_ok = line_color == '16A085'

        if width_ok and color_ok:
            print(f"PASS: Component 5 — Border {line_width/12700:.2f}pt #{line_color} (0.15 pts)")
            total_score += 0.15
        elif color_ok:
            print(f"PARTIAL: Component 5 — Color correct #{line_color}, width {line_width/12700:.2f}pt "
                  f"(expected 1.5pt) (0.07 pts)")
            total_score += 0.07
        elif width_ok:
            print(f"PARTIAL: Component 5 — Width correct {line_width/12700:.2f}pt, color #{line_color} "
                  f"(expected #16A085) (0.07 pts)")
            total_score += 0.07
        else:
            print(f"FAIL: Component 5 — Border width={line_width}, color={line_color}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Text 'A' inside in 24pt bold white (0.15 points)
    try:
        text_found = False
        text_correct = False
        bold_correct = False
        size_correct = False
        color_correct = False

        if triangle.has_text_frame:
            for para in triangle.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip() == 'A':
                        text_found = True
                        text_correct = True
                        # Bold check (None treated as False)
                        bold_correct = run.font.bold is True
                        # Size check: 24pt = 304800 EMU
                        if run.font.size is not None:
                            size_correct = abs(run.font.size - 304800) / 304800 <= 0.05
                        # Color check: white = FFFFFF
                        try:
                            font_rgb = str(run.font.color.rgb).upper()
                            color_correct = font_rgb == 'FFFFFF'
                        except Exception:
                            pass
                        break
                if text_found:
                    break

        sub_checks = sum([text_correct, bold_correct, size_correct, color_correct])
        if sub_checks == 4:
            print(f"PASS: Component 6 — Text 'A', 24pt, bold, white (0.15 pts)")
            total_score += 0.15
        elif text_correct and sub_checks >= 2:
            partial = round(0.15 * sub_checks / 4, 2)
            print(f"PARTIAL: Component 6 — text={text_correct}, bold={bold_correct}, "
                  f"size={size_correct}, color={color_correct} ({partial} pts)")
            total_score += partial
        elif text_correct:
            print(f"PARTIAL: Component 6 — Text 'A' found but properties wrong: "
                  f"bold={bold_correct}, size={size_correct}, color={color_correct} (0.04 pts)")
            total_score += 0.04
        else:
            print(f"FAIL: Component 6 — Text 'A' not found in triangle")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
