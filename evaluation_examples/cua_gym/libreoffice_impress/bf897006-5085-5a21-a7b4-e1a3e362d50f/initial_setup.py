"""
Initial Setup: Create annual report presentation with Fade transitions (3s)
Task ID: impress_tm_019
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import zipfile
import shutil
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_019'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
TEMP_OUTPUT = f'{WORKDIR}/{TASK_ID}_tmp.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide content for a 15-slide annual report
    slide_content = [
        {
            "layout": 0,
            "title": "Meridian Technologies",
            "subtitle": "Annual Report 2025\nPrepared by the Office of the CFO",
        },
        {
            "layout": 1,
            "title": "Executive Summary",
            "body": "Meridian Technologies achieved record revenue of $2.4B in fiscal year 2025, representing a 18% year-over-year increase. Our cloud services division drove growth with 34% expansion, while our enterprise solutions maintained steady margins at 42%. Key strategic acquisitions in Q2 and Q3 strengthened our AI capabilities portfolio.",
        },
        {
            "layout": 1,
            "title": "Revenue Highlights",
            "body": "Total Revenue: $2.41 Billion (+18% YoY)\nCloud Services: $892M (+34% YoY)\nEnterprise Solutions: $764M (+12% YoY)\nProfessional Services: $421M (+8% YoY)\nLicensing & Support: $333M (+5% YoY)",
        },
        {
            "layout": 1,
            "title": "Q1 Performance",
            "body": "Revenue: $547M (vs $468M prior year)\nNew enterprise contracts: 47\nCustomer retention rate: 96.2%\nNotable wins: Pacific Health Systems ($12M), Vertex Financial ($8.5M)\nCloud migration projects initiated: 23",
        },
        {
            "layout": 1,
            "title": "Q2 Performance",
            "body": "Revenue: $589M (vs $502M prior year)\nAcquisition of DataStream Analytics completed ($145M)\nLaunched Meridian AI Assistant platform\nExpanded APAC operations to 3 new markets\nEmployee headcount: 8,450 (+320 from Q1)",
        },
        {
            "layout": 1,
            "title": "Q3 Performance",
            "body": "Revenue: $621M (vs $519M prior year)\nAcquisition of CloudSecure Inc. ($89M)\nMeridian AI Assistant reached 50,000 active users\nOpened new R&D center in Austin, TX\nFiled 34 new patents in cloud infrastructure",
        },
        {
            "layout": 1,
            "title": "Q4 Performance",
            "body": "Revenue: $653M (vs $548M prior year)\nAnnual recurring revenue surpassed $1.8B\nLaunched Meridian Edge Computing Suite\nSigned strategic partnership with GlobalTech Corp\nCustomer NPS score: 72 (industry avg: 54)",
        },
        {
            "layout": 1,
            "title": "Regional Breakdown",
            "body": "North America: $1.42B (59% of total)\nEurope: $578M (24% of total)\nAsia-Pacific: $289M (12% of total)\nRest of World: $123M (5% of total)\nFastest growing region: APAC (+41% YoY)",
        },
        {
            "layout": 1,
            "title": "Product Innovation",
            "body": "Released 4 major product updates\nMeridian Cloud Platform v5.0 with AI-native architecture\nZero-trust security framework integrated across all products\n12 new API integrations with major SaaS platforms\nR&D investment: $386M (16% of revenue)",
        },
        {
            "layout": 1,
            "title": "Customer Success Stories",
            "body": "Pacific Health Systems: Reduced infrastructure costs by 45%\nVertex Financial: Achieved 99.99% uptime on Meridian Cloud\nAutoNova Manufacturing: Streamlined supply chain with AI analytics\nEduConnect: Scaled to 2M concurrent users during peak enrollment",
        },
        {
            "layout": 1,
            "title": "Workforce & Culture",
            "body": "Total employees: 9,200 (up from 7,800)\nDiversity index improved to 0.78 (target: 0.80)\nEmployee satisfaction score: 4.2/5.0\nVoluntary turnover: 8.3% (industry avg: 13.1%)\nLaunched Meridian Academy with 200+ learning paths",
        },
        {
            "layout": 1,
            "title": "Sustainability Initiatives",
            "body": "Carbon neutral operations achieved in Q3 2025\nData centers powered by 100% renewable energy\nReduced water consumption by 28% across facilities\nE-waste recycling rate: 94%\nPlanted 50,000 trees through reforestation partnership",
        },
        {
            "layout": 1,
            "title": "Financial Outlook FY2026",
            "body": "Projected revenue: $2.85B - $2.95B (+18-22%)\nCloud services target: $1.2B\nPlanned R&D investment: $450M+\nExpected operating margin: 22-24%\nCapital expenditure budget: $180M for infrastructure expansion",
        },
        {
            "layout": 1,
            "title": "Strategic Priorities",
            "body": "1. Accelerate AI integration across product portfolio\n2. Expand enterprise market share in EMEA and APAC\n3. Achieve $2B annual recurring revenue milestone\n4. Launch next-generation edge computing platform\n5. Strengthen cybersecurity offerings through M&A",
        },
        {
            "layout": 1,
            "title": "Thank You",
            "body": "Meridian Technologies\nInvestor Relations: ir@meridiantech.com\nMedia Inquiries: press@meridiantech.com\nwww.meridiantech.com/investors\n\nForward-looking statements in this presentation are subject to risks and uncertainties.",
        },
    ]

    for i, content in enumerate(slide_content):
        layout_idx = content["layout"]
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if slide.shapes.title:
            slide.shapes.title.text = content["title"]
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(36) if layout_idx == 0 else Pt(28)
                run.font.bold = True

        if layout_idx == 0 and 1 in slide.placeholders:
            slide.placeholders[1].text = content.get("subtitle", "")
        elif layout_idx == 1 and 1 in slide.placeholders:
            slide.placeholders[1].text = content.get("body", "")

    # Save base pptx first
    prs.save(TEMP_OUTPUT)

    # Now inject Fade transitions (3.0s) into each slide via XML
    inject_transitions(TEMP_OUTPUT, OUTPUT, 'fade', 3000)

    # Clean up temp
    if os.path.exists(TEMP_OUTPUT):
        os.remove(TEMP_OUTPUT)

    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


def inject_transitions(input_path, output_path, transition_type, duration_ms):
    """Inject transitions into all slides of a pptx file.

    transition_type: 'fade', 'dissolve', 'push', etc.
    duration_ms: duration in milliseconds (e.g. 3000 for 3 seconds)
    """
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    # Register namespaces to preserve them in output
    namespaces = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': P_NS,
    }
    for prefix, uri in namespaces.items():
        ET.register_namespace(prefix, uri)

    temp_dir = f'{WORKDIR}/_pptx_temp'
    if os.path.exists(temp_dir):
        shutil.rmtree(temp_dir)

    # Extract pptx (it's a zip)
    with zipfile.ZipFile(input_path, 'r') as zf:
        zf.extractall(temp_dir)

    # Find all slide XML files
    slides_dir = os.path.join(temp_dir, 'ppt', 'slides')
    slide_files = sorted(
        [f for f in os.listdir(slides_dir) if f.startswith('slide') and f.endswith('.xml')],
        key=lambda x: int(x.replace('slide', '').replace('.xml', ''))
    )

    for slide_file in slide_files:
        slide_path = os.path.join(slides_dir, slide_file)
        tree = ET.parse(slide_path)
        root = tree.getroot()

        # Remove existing transition if present
        for tr in root.findall(f'{{{P_NS}}}transition'):
            root.remove(tr)

        # Create transition element
        transition_elem = ET.SubElement(root, f'{{{P_NS}}}transition')
        transition_elem.set('spd', 'slow')
        transition_elem.set('dur', str(duration_ms))

        # Add transition type child element
        ET.SubElement(transition_elem, f'{{{P_NS}}}{transition_type}')

        # Write modified XML back to disk
        tree.write(slide_path, xml_declaration=True, encoding='UTF-8')

    # Repackage as zip
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf_out:
        for root_dir, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root_dir, file)
                arcname = os.path.relpath(file_path, temp_dir)
                zf_out.write(file_path, arcname)

    # Clean up
    shutil.rmtree(temp_dir)


create_initial()
