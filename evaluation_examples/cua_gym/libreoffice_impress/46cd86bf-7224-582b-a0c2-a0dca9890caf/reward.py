"""
Reward Script: UX Research Findings Presentation
Task ID: impress_wf_086
Domain: libreoffice_impress
Scoring: 10 components verifying slide count, content structure, shapes, and visual elements
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_086'


def persist_app_state(domain):
    """Save any unsaved LibreOffice work before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            import time
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_all_text(slide):
    """Recursively extract all text from a slide's shapes."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    texts.append(para.text.strip())
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                if hasattr(sub, 'text_frame') and sub.has_text_frame:
                    for para in sub.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text.strip())
    return texts


def count_shapes_by_type(slide, shape_type):
    """Count shapes of a specific MSO type on a slide."""
    return sum(1 for s in slide.shapes if s.shape_type == shape_type)


def count_auto_shapes_by_name_prefix(slide, prefix):
    """Count auto shapes whose name starts with a given prefix."""
    return sum(1 for s in slide.shapes
               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.name.startswith(prefix))


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = prs.slides
    num_slides = len(slides)

    # Component 1: Exactly 10 slides (0.15 points)
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — Slide count is 10 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 2 slides to check further
    if num_slides < 2:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title text contains 'UX Research Report' (0.10 points)
    try:
        slide1_texts = get_all_text(slides[0])
        slide1_combined = " ".join(slide1_texts).lower()
        if "ux research report" in slide1_combined:
            print(f"PASS: Component 2 — Slide 1 title contains 'UX Research Report' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Slide 1 text does not contain 'UX Research Report'. Found: {slide1_texts[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has 4 method cards (Rounded Rectangles with method names) (0.10 points)
    try:
        if num_slides >= 2:
            slide2 = slides[1]
            slide2_texts = get_all_text(slide2)
            slide2_combined = " ".join(slide2_texts).lower()
            methods_found = sum(1 for m in ["surveys", "interviews", "usability test", "analytics"]
                               if m in slide2_combined)
            rounded_rects = count_auto_shapes_by_name_prefix(slide2, "Rounded Rectangle")
            # Also count any auto shapes that aren't text boxes
            auto_shapes = sum(1 for s in slide2.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
            if methods_found >= 4 and auto_shapes >= 4:
                print(f"PASS: Component 3 — Slide 2 has 4 method cards ({methods_found} methods, {auto_shapes} auto shapes) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Slide 2: {methods_found}/4 methods found, {auto_shapes} auto shapes")
        else:
            print(f"FAIL: Component 3 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has 2 pie chart shapes (side-by-side) (0.10 points)
    try:
        if num_slides >= 3:
            slide3 = slides[2]
            # Pie shapes are AUTO_SHAPE with name starting with "Pie" or "Oval"
            pie_count = sum(1 for s in slide3.shapes
                           if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.name.startswith("Pie"))
            oval_count = sum(1 for s in slide3.shapes
                            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.name.startswith("Oval"))
            # The golden has 2 Pie shapes and 2 Oval shapes on slide 3
            if pie_count >= 2 or (oval_count >= 2 and pie_count >= 1):
                print(f"PASS: Component 4 — Slide 3 has {pie_count} pie shapes, {oval_count} ovals (0.10 pts)")
                total_score += 0.10
            else:
                # Fallback: check for at least 2 circular-ish shapes (could be pies done differently)
                slide3_texts = get_all_text(slide3)
                slide3_combined = " ".join(slide3_texts).lower()
                has_demographics = "demographic" in slide3_combined or "age" in slide3_combined or "device" in slide3_combined
                if has_demographics and (pie_count >= 1 or oval_count >= 2):
                    print(f"PASS: Component 4 — Slide 3 has demographics with charts ({pie_count} pies, {oval_count} ovals) (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 4 — Slide 3: {pie_count} pie shapes, {oval_count} ovals")
        else:
            print(f"FAIL: Component 4 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has 5 numbered insight cards (0.10 points)
    try:
        if num_slides >= 4:
            slide4 = slides[3]
            slide4_texts = get_all_text(slide4)
            # Check for numbered circles (1-5)
            numbers_found = set()
            for txt in slide4_texts:
                txt_stripped = txt.strip()
                if txt_stripped in ["1", "2", "3", "4", "5"]:
                    numbers_found.add(txt_stripped)
            # Count rounded rectangles (card shapes)
            rounded_rects = sum(1 for s in slide4.shapes
                                if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                                and "Rounded Rectangle" in s.name)
            if len(numbers_found) >= 5 and rounded_rects >= 5:
                print(f"PASS: Component 5 — Slide 4 has 5 numbered cards ({len(numbers_found)} numbers, {rounded_rects} cards) (0.10 pts)")
                total_score += 0.10
            elif len(numbers_found) >= 3 and rounded_rects >= 3:
                print(f"PARTIAL: Component 5 — Slide 4 has {len(numbers_found)} numbers, {rounded_rects} cards (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 5 — Slide 4: {len(numbers_found)} numbered items, {rounded_rects} rounded rects")
        else:
            print(f"FAIL: Component 5 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 has emotion curve (circles at varying heights + connector lines) (0.10 points)
    try:
        if num_slides >= 5:
            slide5 = slides[4]
            ovals = [s for s in slide5.shapes
                     if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.name.startswith("Oval")]
            lines = [s for s in slide5.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
            # Golden has 6 ovals at varying heights and 5 connector lines
            # Check that ovals have varying top positions (emotion curve)
            if len(ovals) >= 4 and len(lines) >= 3:
                tops = [o.top for o in ovals]
                unique_tops = len(set(tops))
                if unique_tops >= 3:
                    print(f"PASS: Component 6 — Slide 5 has emotion curve ({len(ovals)} dots, {len(lines)} connectors, {unique_tops} height levels) (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 6 — Slide 5 ovals don't vary in height: {unique_tops} unique tops")
            else:
                print(f"FAIL: Component 6 — Slide 5: {len(ovals)} ovals, {len(lines)} lines (need >=4 ovals, >=3 lines)")
        else:
            print(f"FAIL: Component 6 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 6 has a rectangle with gradient fill (0.10 points)
    try:
        if num_slides >= 6:
            slide6 = slides[5]
            gradient_shapes = [s for s in slide6.shapes
                               if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                               and s.fill.type is not None and s.fill.type == 3]
            if len(gradient_shapes) >= 1:
                print(f"  Found gradient fill on {gradient_shapes[0].name}")
                print(f"PASS: Component 7 — Slide 6 has rectangle with gradient fill (0.10 pts)")
                total_score += 0.10
            else:
                # Check for at least a large rectangle (heatmap placeholder)
                rects = [s for s in slide6.shapes
                         if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "Rectangle" in s.name]
                if len(rects) >= 1:
                    print(f"PARTIAL: Component 7 — Slide 6 has rectangle but no gradient fill (0.05 pts)")
                    total_score += 0.05
                else:
                    print(f"FAIL: Component 7 — Slide 6: no heatmap rectangle found")
        else:
            print(f"FAIL: Component 7 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 7 has horizontal bar chart (multiple rectangles with % labels) (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = slides[6]
            rects = [s for s in slide7.shapes
                     if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "Rectangle" in s.name]
            texts = get_all_text(slide7)
            pct_labels = [t for t in texts if "%" in t]
            if len(rects) >= 5 and len(pct_labels) >= 5:
                print(f"PASS: Component 8 — Slide 7 has bar chart ({len(rects)} bars, {len(pct_labels)} % labels) (0.10 pts)")
                total_score += 0.10
            elif len(rects) >= 3 and len(pct_labels) >= 3:
                print(f"PARTIAL: Component 8 — Slide 7 has partial bar chart ({len(rects)} bars, {len(pct_labels)} labels) (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 — Slide 7: {len(rects)} rectangles, {len(pct_labels)} % labels")
        else:
            print(f"FAIL: Component 8 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 8 has SUS gauge (arc/semi-circle shape + line/needle) (0.10 points)
    try:
        if num_slides >= 8:
            slide8 = slides[7]
            arcs = [s for s in slide8.shapes
                    if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "Arc" in s.name]
            lines = [s for s in slide8.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
            texts = get_all_text(slide8)
            texts_combined = " ".join(texts).lower()
            has_sus_ref = "sus" in texts_combined or "usability" in texts_combined or "72" in texts_combined or "score" in texts_combined
            if len(arcs) >= 1 and len(lines) >= 1 and has_sus_ref:
                print(f"PASS: Component 9 — Slide 8 has SUS gauge ({len(arcs)} arcs, {len(lines)} needle lines) (0.10 pts)")
                total_score += 0.10
            elif has_sus_ref and (len(arcs) >= 1 or len(lines) >= 1):
                print(f"PARTIAL: Component 9 — Slide 8 partial gauge (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 9 — Slide 8: {len(arcs)} arcs, {len(lines)} lines, SUS ref={has_sus_ref}")
        else:
            print(f"FAIL: Component 9 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Slide 9 has a table (Impact vs Effort matrix, at least 2x2) (0.05 points)
    try:
        if num_slides >= 9:
            slide9 = slides[8]
            qualifying_tables = [s for s in slide9.shapes
                                 if s.shape_type == MSO_SHAPE_TYPE.TABLE
                                 and len(s.table.rows) >= 2
                                 and len(s.table.columns) >= 2]
            if len(qualifying_tables) >= 1:
                t = qualifying_tables[0].table
                print(f"  Found {len(t.rows)}x{len(t.columns)} table on slide 9")
                print(f"PASS: Component 10 — Slide 9 has Impact/Effort matrix table (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 10 — Slide 9: no qualifying table found")
        else:
            print(f"FAIL: Component 10 — Not enough slides")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/UX_Research.pptx'
if not os.path.exists(file_path):
    # Also check Desktop
    desktop_path = f'{WORKDIR}/Desktop/UX_Research.pptx'
    if os.path.exists(desktop_path):
        file_path = desktop_path
    else:
        print(f"File not found: {file_path} or {desktop_path}")
        print("REWARD: 0.0")
        exit()

verify_task(file_path)
