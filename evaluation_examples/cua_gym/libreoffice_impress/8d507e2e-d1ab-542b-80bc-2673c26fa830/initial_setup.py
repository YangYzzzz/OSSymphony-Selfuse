"""
Initial Setup: Convert bullet list to SmartArt-style visual layout
Task ID: impress_rp_041
Domain: libreoffice_impress

Creates a 10-slide presentation with slide 4 containing a title 'Our Priorities'
and 8 bullet points. The agent's task is to convert these bullets into rounded
rectangle shapes.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_041'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Key Initiatives 2026"
    slide1.placeholders[1].text = "Strategic Planning & Roadmap Overview"

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Company Vision & Mission"
    items2 = [
        "Market Analysis & Competitive Landscape",
        "Strategic Priorities Overview",
        "Our Priorities",
        "Implementation Timeline",
        "Resource Allocation",
        "Key Performance Indicators",
        "Q&A Session",
    ]
    for item in items2:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # ---- Slide 3: Vision & Mission ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Vision & Mission"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Our vision is to become the leading technology partner for Fortune 500 enterprises."
    p3 = body3.add_paragraph()
    p3.text = ""
    p3b = body3.add_paragraph()
    p3b.text = "Mission: Deliver innovative, scalable solutions that transform how businesses operate in the digital age."

    # ---- Slide 4: Our Priorities (the target slide) ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title text box
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Our Priorities"
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.runs[0]
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Bullet list text box
    bullet_items = [
        "Customer Experience",
        "Data Analytics",
        "Cloud Migration",
        "Talent Development",
        "Sustainability",
        "Innovation Lab",
        "Global Expansion",
        "Cybersecurity",
    ]
    bullet_box = slide4.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10), Inches(5.0))
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True

    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf_bullets.paragraphs[0]
        else:
            p = tf_bullets.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ---- Slide 5: Market Analysis ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Market Analysis"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Total addressable market: $142B by 2027"
    for txt in [
        "Annual growth rate: 14.2% CAGR",
        "Key segments: Enterprise SaaS, Cloud Infrastructure, AI/ML Services",
        "Competitive advantage: Proprietary platform with 99.97% uptime",
    ]:
        p = body5.add_paragraph()
        p.text = txt

    # ---- Slide 6: Implementation Timeline ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Implementation Timeline"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Q1 2026: Foundation & Infrastructure Setup"
    for txt in [
        "Q2 2026: Pilot Programs & Early Adoption",
        "Q3 2026: Full Rollout & Integration",
        "Q4 2026: Optimization & Performance Review",
    ]:
        p = body6.add_paragraph()
        p.text = txt

    # ---- Slide 7: Resource Allocation ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Budget & Resources"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Total Budget: $24.5M allocated across 8 initiatives"
    for txt in [
        "Engineering: 45% ($11.0M)",
        "Operations: 25% ($6.1M)",
        "Marketing & Sales: 20% ($4.9M)",
        "Research & Innovation: 10% ($2.5M)",
    ]:
        p = body7.add_paragraph()
        p.text = txt

    # ---- Slide 8: KPIs ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Key Performance Indicators"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Customer satisfaction score: Target 92%+"
    for txt in [
        "Revenue growth: 18% year-over-year",
        "Employee retention rate: >90%",
        "Cloud migration completion: 85% by Q4",
        "Innovation pipeline: 12+ new features per quarter",
    ]:
        p = body8.add_paragraph()
        p.text = txt

    # ---- Slide 9: Risk Assessment ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Risk Assessment"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Talent acquisition in competitive AI/ML market"
    for txt in [
        "Supply chain dependencies for cloud infrastructure",
        "Regulatory compliance across global markets",
        "Cybersecurity threat landscape evolution",
    ]:
        p = body9.add_paragraph()
        p.text = txt

    # ---- Slide 10: Next Steps ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Next Steps & Q&A"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Review and finalize initiative ownership by March 28"
    for txt in [
        "Schedule kick-off meetings for each workstream",
        "Establish monthly progress review cadence",
        "Questions and discussion",
    ]:
        p = body10.add_paragraph()
        p.text = txt

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
