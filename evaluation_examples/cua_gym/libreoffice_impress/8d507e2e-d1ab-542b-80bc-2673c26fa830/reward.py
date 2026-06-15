"""
Reward Script: Convert bullet list to SmartArt-style rounded rectangles
Task ID: impress_rp_041
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 8 rounded rectangle shapes exist on slide 4
  Component 2 (0.25): Each rectangle contains the correct bullet text in order
  Component 3 (0.25): Alternating fill colors (#3498DB odd, #2ECC71 even)
  Component 4 (0.20): White bold 12pt text in each rectangle
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_041'

# Expected bullet texts in order (top-left to bottom-right, row by row)
EXPECTED_TEXTS = [
    'Customer Experience',
    'Data Analytics',
    'Cloud Migration',
    'Talent Development',
    'Sustainability',
    'Innovation Lab',
    'Global Expansion',
    'Cybersecurity',
]

# Expected alternating fill colors (1-indexed: odd=#3498DB, even=#2ECC71)
ODD_COLOR = '3498DB'
EVEN_COLOR = '2ECC71'


def get_rounded_rects(slide):
    """Find all rounded rectangle auto shapes on a slide, sorted by position (top-to-bottom, left-to-right)."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if geometry is roundRect
            sp_el = shape._element
            prst_el = sp_el.find('.//' + qn('a:prstGeom'))
            if prst_el is not None and prst_el.get('prst') == 'roundRect':
                rects.append(shape)
    # Sort by top (row), then left (column) to get reading order
    rects.sort(key=lambda s: (s.top, s.left))
    return rects


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Get all rounded rectangles on slide 4
    rects = get_rounded_rects(slide)

    # Component 1: 8 rounded rectangle shapes exist on slide 4 (0.30 points)
    try:
        num_rects = len(rects)
        if num_rects == 8:
            print(f"PASS: Component 1 -- Found exactly 8 rounded rectangles (0.30 pts)")
            total_score += 0.30
        elif num_rects >= 6:
            partial = 0.15
            print(f"PARTIAL: Component 1 -- Found {num_rects} rounded rectangles, expected 8 ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Found {num_rects} rounded rectangles, expected 8")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if num_rects == 0:
        # No rectangles at all, cannot check further
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Each rectangle contains the correct bullet text (0.25 points)
    try:
        correct_text_count = 0
        for i, rect in enumerate(rects):
            if i >= len(EXPECTED_TEXTS):
                break
            actual_text = rect.text.strip() if hasattr(rect, 'text') else ''
            expected_text = EXPECTED_TEXTS[i]
            if actual_text == expected_text:
                correct_text_count += 1
            else:
                print(f"  DETAIL: Rect {i+1} text mismatch: expected '{expected_text}', found '{actual_text}'")

        if correct_text_count == 8:
            print(f"PASS: Component 2 -- All 8 rectangles have correct text (0.25 pts)")
            total_score += 0.25
        elif correct_text_count >= 4:
            partial = round(0.25 * correct_text_count / 8, 2)
            print(f"PARTIAL: Component 2 -- {correct_text_count}/8 correct texts ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- Only {correct_text_count}/8 correct texts")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Alternating fill colors (0.25 points)
    # Odd items (1,3,5,7 = indices 0,2,4,6) -> #3498DB
    # Even items (2,4,6,8 = indices 1,3,5,7) -> #2ECC71
    try:
        correct_color_count = 0
        for i, rect in enumerate(rects):
            if i >= 8:
                break
            expected_color = ODD_COLOR if i % 2 == 0 else EVEN_COLOR
            try:
                fill = rect.fill
                if fill.type is not None:
                    actual_color = str(fill.fore_color.rgb)
                else:
                    actual_color = 'no_fill'
            except Exception:
                actual_color = 'error'

            if actual_color.upper() == expected_color.upper():
                correct_color_count += 1
            else:
                print(f"  DETAIL: Rect {i+1} color mismatch: expected #{expected_color}, found #{actual_color}")

        if correct_color_count == 8:
            print(f"PASS: Component 3 -- All 8 rectangles have correct alternating fill colors (0.25 pts)")
            total_score += 0.25
        elif correct_color_count >= 4:
            partial = round(0.25 * correct_color_count / 8, 2)
            print(f"PARTIAL: Component 3 -- {correct_color_count}/8 correct colors ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Only {correct_color_count}/8 correct colors")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: White bold 12pt text formatting (0.20 points)
    try:
        correct_format_count = 0
        for i, rect in enumerate(rects):
            if i >= 8:
                break
            if not rect.has_text_frame:
                print(f"  DETAIL: Rect {i+1} has no text frame")
                continue

            tf = rect.text_frame
            # Check all non-empty runs
            all_runs_ok = True
            found_runs = False
            for para in tf.paragraphs:
                for run in para.runs:
                    if not (run.text or '').strip():
                        continue
                    found_runs = True

                    # Check bold
                    is_bold = run.font.bold is True
                    if not is_bold:
                        all_runs_ok = False
                        print(f"  DETAIL: Rect {i+1} run not bold (bold={run.font.bold})")

                    # Check white color (FFFFFF)
                    try:
                        if run.font.color.type is not None:
                            color_rgb = str(run.font.color.rgb).upper()
                            if color_rgb != 'FFFFFF':
                                all_runs_ok = False
                                print(f"  DETAIL: Rect {i+1} run color={color_rgb}, expected FFFFFF")
                        else:
                            all_runs_ok = False
                            print(f"  DETAIL: Rect {i+1} run color type is None (inherited)")
                    except Exception:
                        all_runs_ok = False

                    # Check 12pt (152400 EMU)
                    if run.font.size is not None:
                        size_pt = run.font.size / 12700  # EMU to pt
                        if abs(size_pt - 12.0) > 0.5:
                            all_runs_ok = False
                            print(f"  DETAIL: Rect {i+1} run size={size_pt}pt, expected 12pt")
                    else:
                        all_runs_ok = False
                        print(f"  DETAIL: Rect {i+1} run size is None (inherited)")

            if found_runs and all_runs_ok:
                correct_format_count += 1

        if correct_format_count == 8:
            print(f"PASS: Component 4 -- All 8 rectangles have white bold 12pt text (0.20 pts)")
            total_score += 0.20
        elif correct_format_count >= 4:
            partial = round(0.20 * correct_format_count / 8, 2)
            print(f"PARTIAL: Component 4 -- {correct_format_count}/8 correctly formatted ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Only {correct_format_count}/8 correctly formatted")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
