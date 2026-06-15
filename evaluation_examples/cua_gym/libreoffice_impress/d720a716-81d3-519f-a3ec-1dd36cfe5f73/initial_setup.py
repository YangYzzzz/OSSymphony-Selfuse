"""
Initial Setup: Photo gallery presentation with 12 slides
Task ID: impress_rp_018
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_018'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Helper to add a title-only slide
    def add_title_slide(title_text, subtitle_text=None):
        layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        if subtitle_text and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle_text
        return slide

    def add_content_slide(title_text, body_lines=None):
        layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        if body_lines and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.paragraphs[0].text = body_lines[0]
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.level = 0
        return slide

    def add_blank_with_title(title_text):
        layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(layout)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        return slide

    # Slide 1: Title slide
    add_title_slide(
        "Annual Company Retreat 2025",
        "Lakewood Resort, Colorado | June 14-16, 2025"
    )

    # Slide 2: Agenda
    add_content_slide("Weekend Agenda", [
        "Friday: Arrival & Welcome Dinner",
        "Saturday: Team Building Activities & Workshops",
        "Saturday Evening: Outdoor BBQ & Bonfire",
        "Sunday: Awards Ceremony & Departure",
    ])

    # Slide 3: Venue Overview
    add_content_slide("Venue Overview", [
        "Lakewood Mountain Resort — 150 acres of pristine nature",
        "Conference center with 3 breakout rooms",
        "Outdoor amphitheater for evening events",
        "On-site dining for up to 200 guests",
    ])

    # Slide 4: Team Building Activities
    add_content_slide("Team Building Activities", [
        "Kayaking Challenge on Mirror Lake",
        "Wilderness Navigation Course",
        "Collaborative Mural Painting",
        "Escape Room: The Mountain Mystery",
        "Cooking Competition: Camp Chef Showdown",
    ])

    # Slide 5: Workshop Schedule
    add_content_slide("Workshop Schedule", [
        "9:00 AM — Leadership in Action (Grand Hall)",
        "10:30 AM — Creative Problem Solving (Room B)",
        "1:00 PM — Cross-Team Collaboration (Outdoor Pavilion)",
        "3:00 PM — Innovation Sprint (Conference Room A)",
    ])

    # Slide 6: Awards & Recognition
    add_content_slide("Awards & Recognition", [
        "Employee of the Year: Outstanding Contribution",
        "Team Excellence Award: Best Collaboration",
        "Innovation Champion: Creative Solutions",
        "Community Impact Award: Volunteer Leadership",
    ])

    # Slide 7: Logistics & Travel
    add_content_slide("Logistics & Travel", [
        "Shuttle service from Denver International Airport",
        "Check-in begins Friday at 2:00 PM",
        "Dress code: Smart casual (outdoor-appropriate)",
        "Contact: events@company.com | (303) 555-0147",
    ])

    # Slide 8: Event Gallery — TITLE ONLY, NO PLACEHOLDERS
    slide8 = add_blank_with_title("Event Gallery")
    # Intentionally empty — the task is to add photo placeholders here

    # Slide 9: Participant Feedback
    add_content_slide("Participant Feedback", [
        '"Best retreat we have ever had!" — Sarah Chen, Engineering',
        '"The team building activities were incredible." — Marcus Rivera, Sales',
        '"I loved the outdoor workshops." — Priya Sharma, Product Design',
        '"Looking forward to next year already!" — James O\'Brien, Finance',
    ])

    # Slide 10: Budget Summary
    add_content_slide("Budget Summary", [
        "Venue & Accommodation: $28,500",
        "Catering & Dining: $12,300",
        "Activities & Equipment: $8,750",
        "Transportation: $4,200",
        "Total: $53,750 (under budget by $1,250)",
    ])

    # Slide 11: Next Steps
    add_content_slide("Next Steps", [
        "Distribute post-event survey by June 20",
        "Compile photo gallery for company newsletter",
        "Schedule planning meeting for 2026 retreat",
        "Submit final expense reports by June 30",
    ])

    # Slide 12: Thank You
    add_title_slide(
        "Thank You!",
        "See you at the 2026 retreat!"
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
