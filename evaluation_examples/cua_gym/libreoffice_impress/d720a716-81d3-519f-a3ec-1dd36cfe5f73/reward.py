"""
Reward Script: Photo gallery slide with 6 rounded rectangle placeholders in 3x2 grid
Task ID: impress_rp_018
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Six rounded rectangle shapes on slide 8
  Component 2 (0.25): Fill #E0E0E0 and border #BDBDBD at 1pt on all rectangles
  Component 3 (0.25): Six caption text boxes with 'Caption' text in ~10pt gray
  Component 4 (0.20): 3x2 grid layout (3 distinct columns, 2 distinct rows)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_018'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu, Pt
        from pptx.dml.color import RGBColor
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # slide 8, 0-indexed

    # Collect rounded rectangles and caption text boxes
    rounded_rects = []
    caption_boxes = []

    for shape in slide.shapes:
        # Rounded rectangles: AUTO_SHAPE type with auto_shape_type == ROUNDED_RECTANGLE (5)
        if shape.shape_type == 1:  # AUTO_SHAPE
            try:
                if shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                    rounded_rects.append(shape)
            except Exception:
                pass
        # Caption text boxes: TEXT_BOX with text containing 'caption' (case-insensitive)
        # Exclude the title text box (which has 'Event Gallery')
        elif shape.shape_type == 17:  # TEXT_BOX
            if shape.has_text_frame:
                text = shape.text_frame.text.strip().lower()
                if 'caption' in text:
                    caption_boxes.append(shape)

    print(f"Found {len(rounded_rects)} rounded rectangles")
    print(f"Found {len(caption_boxes)} caption text boxes")

    # Component 1: Six rounded rectangle shapes on slide 8 (0.30 points)
    try:
        if len(rounded_rects) >= 6:
            print(f"PASS: Component 1 — Found {len(rounded_rects)} rounded rectangles (0.30 pts)")
            total_score += 0.30
        elif len(rounded_rects) >= 3:
            partial = 0.15
            print(f"PARTIAL: Component 1 — Found {len(rounded_rects)}/6 rounded rectangles ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — Found {len(rounded_rects)}/6 rounded rectangles")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Fill #E0E0E0 and border #BDBDBD at ~1pt on all rectangles (0.25 points)
    try:
        if len(rounded_rects) >= 6:
            fill_ok_count = 0
            border_ok_count = 0
            for rect in rounded_rects[:6]:
                # Check fill color
                try:
                    if rect.fill.type == 1:  # SOLID
                        fill_rgb = str(rect.fill.fore_color.rgb).upper()
                        if fill_rgb == 'E0E0E0':
                            fill_ok_count += 1
                except Exception:
                    pass
                # Check border color and width
                try:
                    line = rect.line
                    if line.fill.type == 1:  # SOLID
                        line_rgb = str(line.color.rgb).upper()
                        # 1pt = 12700 EMU, allow tolerance
                        line_width = line.width
                        if line_rgb == 'BDBDBD' and abs(line_width - 12700) <= 2000:
                            border_ok_count += 1
                except Exception:
                    pass

            fill_ratio = fill_ok_count / 6.0
            border_ratio = border_ok_count / 6.0
            combined_ratio = (fill_ratio + border_ratio) / 2.0
            comp2_score = round(0.25 * combined_ratio, 4)

            if comp2_score >= 0.24:
                print(f"PASS: Component 2 — Fill OK: {fill_ok_count}/6, Border OK: {border_ok_count}/6 ({comp2_score} pts)")
            elif comp2_score > 0:
                print(f"PARTIAL: Component 2 — Fill OK: {fill_ok_count}/6, Border OK: {border_ok_count}/6 ({comp2_score} pts)")
            else:
                print(f"FAIL: Component 2 — Fill OK: {fill_ok_count}/6, Border OK: {border_ok_count}/6")
            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 — Not enough rounded rectangles to check fill/border")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Six caption text boxes with 'Caption' text in ~10pt gray font (0.25 points)
    try:
        if len(caption_boxes) >= 6:
            style_ok_count = 0
            for cbox in caption_boxes[:6]:
                try:
                    for p in cbox.text_frame.paragraphs:
                        for r in p.runs:
                            # Check font size: 10pt = 127000 EMU, allow tolerance (8-12pt)
                            size_ok = False
                            if r.font.size is not None:
                                size_pt = r.font.size / 12700  # convert EMU to pt
                                if 8 <= size_pt <= 12:
                                    size_ok = True
                            # Check font color is grayish
                            color_ok = False
                            try:
                                if r.font.color.type is not None:
                                    rgb = r.font.color.rgb
                                    r_val = int(str(rgb)[0:2], 16)
                                    g_val = int(str(rgb)[2:4], 16)
                                    b_val = int(str(rgb)[4:6], 16)
                                    # Gray means R==G==B and in the range ~80-180
                                    if abs(r_val - g_val) <= 20 and abs(g_val - b_val) <= 20 and 60 <= r_val <= 200:
                                        color_ok = True
                            except Exception:
                                pass
                            if size_ok and color_ok:
                                style_ok_count += 1
                                break  # one good run is enough per box
                        else:
                            continue
                        break
                except Exception:
                    pass

            # Score based on how many boxes have correct style
            if len(caption_boxes) >= 6:
                count_ratio = min(len(caption_boxes), 6) / 6.0
                style_ratio = style_ok_count / 6.0
                # 50% for having 6 captions, 50% for correct styling
                comp3_score = round(0.25 * (0.5 * count_ratio + 0.5 * style_ratio), 4)
            else:
                comp3_score = round(0.25 * (len(caption_boxes) / 6.0) * 0.5, 4)

            if comp3_score >= 0.24:
                print(f"PASS: Component 3 — {len(caption_boxes)} captions, {style_ok_count}/6 styled correctly ({comp3_score} pts)")
            elif comp3_score > 0:
                print(f"PARTIAL: Component 3 — {len(caption_boxes)} captions, {style_ok_count}/6 styled correctly ({comp3_score} pts)")
            else:
                print(f"FAIL: Component 3 — {len(caption_boxes)} captions, {style_ok_count}/6 styled correctly")
            total_score += comp3_score
        elif len(caption_boxes) > 0:
            partial = round(0.25 * (len(caption_boxes) / 6.0) * 0.5, 4)
            print(f"PARTIAL: Component 3 — Only {len(caption_boxes)}/6 caption boxes ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No caption text boxes found")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 3x2 grid layout — 3 distinct columns, 2 distinct rows (0.20 points)
    try:
        if len(rounded_rects) >= 6:
            # Collect left positions and top positions
            lefts = sorted(set(rect.left for rect in rounded_rects[:6]))
            tops = sorted(set(rect.top for rect in rounded_rects[:6]))

            # Cluster positions (within 5% tolerance)
            def cluster_positions(positions, tolerance_emu=200000):
                """Cluster positions that are within tolerance of each other."""
                clusters = []
                for pos in sorted(positions):
                    merged = False
                    for cluster in clusters:
                        if abs(pos - cluster[0]) <= tolerance_emu:
                            cluster.append(pos)
                            merged = True
                            break
                    if not merged:
                        clusters.append([pos])
                return clusters

            all_lefts = [rect.left for rect in rounded_rects[:6]]
            all_tops = [rect.top for rect in rounded_rects[:6]]

            col_clusters = cluster_positions(all_lefts)
            row_clusters = cluster_positions(all_tops)

            num_cols = len(col_clusters)
            num_rows = len(row_clusters)

            print(f"  Grid analysis: {num_cols} columns, {num_rows} rows")

            if num_cols == 3 and num_rows == 2:
                print(f"PASS: Component 4 — 3x2 grid layout confirmed (0.20 pts)")
                total_score += 0.20
            elif num_cols >= 2 and num_rows >= 2:
                partial = 0.10
                print(f"PARTIAL: Component 4 — Grid is {num_cols}x{num_rows} instead of 3x2 ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — Grid is {num_cols}x{num_rows}, expected 3x2")
        else:
            print(f"FAIL: Component 4 — Not enough rounded rectangles for grid analysis")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
