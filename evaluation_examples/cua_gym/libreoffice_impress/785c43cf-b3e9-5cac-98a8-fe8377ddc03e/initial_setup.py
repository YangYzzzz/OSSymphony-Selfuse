"""
Initial Setup: Apply Corporate_Standard.otp template to Draft_Presentation.pptx
Task ID: impress_fix_042
Domain: libreoffice_impress

Creates:
  1. /home/user/Templates/Corporate_Standard.otp - A branded template file
  2. /home/user/Draft_Presentation.pptx - A 10-slide presentation with plain white master
Opens Draft_Presentation.pptx in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_042'
TEMPLATES_DIR = f'{WORKDIR}/Templates'
TEMPLATE_FILE = f'{TEMPLATES_DIR}/Corporate_Standard.otp'
DRAFT_FILE = f'{WORKDIR}/Draft_Presentation.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_name="Arial",
                font_size=Pt(18), bold=False, color=None, alignment=None):
    """Helper to add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_name="Arial",
                    font_size=Pt(16), color=None):
    """Helper to add a bullet list textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        # Add bullet character
        pPr = p._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '\u2022'}))
        for run in p.runs:
            run.font.name = font_name
            run.font.size = font_size
            if color:
                run.font.color.rgb = color
    return txBox


def create_template():
    """Create the Corporate_Standard.otp template with branded master slides."""
    os.makedirs(TEMPLATES_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Modify the slide master to have corporate branding
    slide_master = prs.slide_masters[0]

    # Set master background to dark navy
    bg = slide_master.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0A, 0x1F, 0x3C)  # Dark navy

    # Add a title slide (layout 0) with just a placeholder feel
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Set slide background to corporate navy
    sfill = slide.background.fill
    sfill.solid()
    sfill.fore_color.rgb = RGBColor(0x0A, 0x1F, 0x3C)

    if slide.shapes.title:
        slide.shapes.title.text = "Corporate Standard Template"
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.name = "Calibri"
            run.font.size = Pt(36)

    # Save as .otp (it's just a pptx with different extension for LO)
    prs.save(TEMPLATE_FILE)
    print(f'Template created: {TEMPLATE_FILE}')


def create_draft_presentation():
    """Create Draft_Presentation.pptx with 10 slides, plain white master, business content."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 2025 Strategic Review"
    slide1.placeholders[1].text = "Meridian Technologies Inc.\nPrepared by Elena Vasquez, VP Strategy"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Agenda", font_size=Pt(32), bold=True, color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide2, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "Market Performance Overview",
        "Product Launch Metrics: Nexus Pro & Horizon Suite",
        "Regional Expansion Update: APAC & EMEA",
        "Competitive Landscape Analysis",
        "Operational Efficiency Gains",
        "Customer Retention & Satisfaction Trends",
        "Financial Projections for Q4 2025",
        "Strategic Priorities for 2026",
    ], font_size=Pt(18), color=RGBColor(0x33, 0x33, 0x33))

    # --- Slide 3: Market Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Market Performance Overview", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide3, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "Total revenue reached $127.4M, up 18.3% year-over-year",
        "Market share grew from 14.2% to 16.8% across all verticals",
        "Enterprise segment contributed $82.1M (64% of total revenue)",
        "SaaS recurring revenue hit $41.2M monthly run rate",
        "Gross margin improved to 72.4%, up from 68.9% in Q2",
    ], font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 4: Product Launch ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Product Launch Metrics", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    # Add a table for product metrics
    table_shape = slide4.shapes.add_table(4, 4, Inches(0.8), Inches(1.5), Inches(8.2), Inches(2.5))
    table = table_shape.table
    headers = ["Product", "Launch Date", "Users (30d)", "Revenue Impact"]
    data = [
        ["Nexus Pro 3.0", "Jul 15, 2025", "34,200", "$8.7M"],
        ["Horizon Suite", "Aug 22, 2025", "18,900", "$5.2M"],
        ["CloudBridge API", "Sep 5, 2025", "12,400", "$3.1M"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 5: Regional Expansion ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Regional Expansion Update", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide5, Inches(0.5), Inches(1.5), Inches(4.2), Inches(5), [
        "APAC: Tokyo office operational since August",
        "APAC: 47 new enterprise clients signed",
        "APAC: Revenue contribution $18.3M (+42% QoQ)",
    ], font_size=Pt(14), color=RGBColor(0x44, 0x44, 0x44))
    add_bullet_list(slide5, Inches(5.2), Inches(1.5), Inches(4.2), Inches(5), [
        "EMEA: Berlin hub expanded to 85 employees",
        "EMEA: Strategic partnership with SAP finalized",
        "EMEA: Revenue contribution $23.7M (+28% QoQ)",
    ], font_size=Pt(14), color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 6: Competitive Landscape ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Competitive Landscape Analysis", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide6, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "Primary competitor Synthex Corp lost 2.3% market share after product recall",
        "New entrant DataPulse raised $120M Series C, targeting mid-market",
        "Our NPS score (78) leads the industry vs. Synthex (62) and Vortex (55)",
        "Patent portfolio expanded to 142 active patents (up from 118)",
    ], font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 7: Operational Efficiency ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Operational Efficiency Gains", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    table_shape7 = slide7.shapes.add_table(5, 3, Inches(0.8), Inches(1.5), Inches(8.2), Inches(3))
    table7 = table_shape7.table
    headers7 = ["Metric", "Q2 2025", "Q3 2025"]
    data7 = [
        ["Avg. Deploy Time", "4.2 hours", "1.8 hours"],
        ["Incident Response", "23 min", "11 min"],
        ["Infrastructure Cost/User", "$4.80", "$3.20"],
        ["Automation Coverage", "62%", "81%"],
    ]
    for c, h in enumerate(headers7):
        cell = table7.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data7, 1):
        for c, val in enumerate(row):
            table7.cell(r, c).text = val

    # --- Slide 8: Customer Retention ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Customer Retention & Satisfaction", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide8, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "Annual retention rate improved to 94.7% (up from 91.2%)",
        "Enterprise tier retention at 98.1%, highest in company history",
        "Average support ticket resolution dropped to 2.4 hours",
        "Customer health score average: 8.6/10 across 1,247 accounts",
        "Expansion revenue from existing accounts: $19.3M (15.1% of total)",
    ], font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 9: Financial Projections ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide9, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Q4 2025 Financial Projections", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide9, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "Projected revenue: $142.8M (conservative) to $151.3M (optimistic)",
        "Expected EBITDA margin: 28-31%",
        "Planned R&D investment: $22.5M for Nexus 4.0 platform",
        "Hiring targets: 120 new positions across engineering and sales",
        "Capital expenditure: $8.4M for new data center in Singapore",
    ], font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

    # --- Slide 10: Strategic Priorities ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide10, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Strategic Priorities for 2026", font_size=Pt(28), bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_bullet_list(slide10, Inches(0.8), Inches(1.5), Inches(8), Inches(5), [
        "1. Achieve $600M annual revenue run rate by end of 2026",
        "2. Launch AI-powered analytics module (Project Athena)",
        "3. Expand LATAM presence with Sao Paulo regional office",
        "4. Pursue strategic acquisition in cybersecurity vertical",
        "5. Achieve SOC 2 Type II and ISO 27001 certifications",
        "6. Build partner ecosystem with 50+ certified integrators",
    ], font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

    prs.save(DRAFT_FILE)
    print(f'Draft presentation created: {DRAFT_FILE}')


def main():
    create_template()
    create_draft_presentation()

    # Launch LibreOffice Impress with the draft presentation
    launch_gui(f'libreoffice --impress "{DRAFT_FILE}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
