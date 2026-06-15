"""
Reward Script: Apply Corporate_Standard.otp template to Draft_Presentation.pptx
Task ID: impress_fix_042
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): Background changed to solid fill on all slides
  Component 2 (0.30): Font changed from Arial to Calibri
  Component 3 (0.20): Font colors changed to template theme (FFFFFF titles, CCD6E0 body)
  Component 4 (0.20): Text content preserved (no content loss)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_042'
FILE_PATH = os.path.join(WORKDIR, 'Draft_Presentation.pptx')

# Known initial-state text content for preservation check
# These are distinctive text snippets from each slide that must survive the template application
EXPECTED_TEXTS = [
    "Q3 2025 Strategic Review",
    "Meridian Technologies Inc.",
    "Agenda",
    "Market Performance Overview",
    "Product Launch Metrics",
    "Regional Expansion Update",
    "Competitive Landscape Analysis",
    "Operational Efficiency Gains",
    "Customer Retention & Satisfaction",
    "Q4 2025 Financial Projections",
    "Strategic Priorities for 2026",
    "$127.4M",
    "Elena Vasquez",
    "Nexus Pro",
    "APAC",
    "EMEA",
    "Synthex Corp",
    "$600M",
]


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have 10 slides
    num_slides = len(prs.slides)
    if num_slides != 10:
        print(f"PRECONDITION FAIL: Expected 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Background changed to solid fill on all slides (0.30 points)
    # Initial state has BACKGROUND(5) = inherited; golden has SOLID(1) = explicit solid fill
    try:
        solid_bg_count = 0
        for i, slide in enumerate(prs.slides):
            fill = slide.background.fill
            if fill.type == 1:  # SOLID fill
                solid_bg_count += 1
            else:
                print(f"  Slide {i+1}: bg fill type = {fill.type} (expected SOLID/1)")

        if solid_bg_count == 10:
            print(f"PASS: Component 1 - All 10 slides have solid background fill (0.30 pts)")
            total_score += 0.30
        elif solid_bg_count >= 5:
            partial = 0.15
            print(f"PARTIAL: Component 1 - {solid_bg_count}/10 slides have solid bg ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 - Only {solid_bg_count}/10 slides have solid background")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Font changed from Arial to Calibri across all text runs (0.30 points)
    # Initial state uses Arial; golden uses Calibri
    try:
        total_named_runs = 0
        calibri_runs = 0
        non_calibri_fonts = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():  # only count non-empty runs
                                if run.font.name is not None:
                                    total_named_runs += 1
                                    if run.font.name == "Calibri":
                                        calibri_runs += 1
                                    else:
                                        non_calibri_fonts.add(run.font.name)

        if total_named_runs > 0:
            calibri_ratio = calibri_runs / total_named_runs
            if calibri_ratio >= 0.9:
                print(f"PASS: Component 2 - {calibri_runs}/{total_named_runs} runs use Calibri (0.30 pts)")
                total_score += 0.30
            elif calibri_ratio >= 0.5:
                partial = 0.15
                print(f"PARTIAL: Component 2 - {calibri_runs}/{total_named_runs} runs use Calibri ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 2 - Only {calibri_runs}/{total_named_runs} runs use Calibri. Found: {non_calibri_fonts}")
        else:
            print(f"FAIL: Component 2 - No runs with explicit font names found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Font colors changed to template theme (0.20 points)
    # Initial: titles 333333, body 444444; Golden: titles FFFFFF, body CCD6E0
    # We check that text runs no longer use the old dark colors (333333, 444444)
    # and instead use the new light colors (FFFFFF, CCD6E0)
    try:
        old_color_count = 0
        new_color_count = 0
        total_colored_runs = 0

        old_colors = {"333333", "444444"}
        new_colors = {"FFFFFF", "CCD6E0"}

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                try:
                                    if run.font.color.type is not None:
                                        rgb_str = str(run.font.color.rgb)
                                        total_colored_runs += 1
                                        if rgb_str in old_colors:
                                            old_color_count += 1
                                        elif rgb_str in new_colors:
                                            new_color_count += 1
                                except Exception:
                                    pass

        if total_colored_runs > 0 and new_color_count > 0 and old_color_count == 0:
            print(f"PASS: Component 3 - {new_color_count}/{total_colored_runs} runs use new theme colors, 0 use old colors (0.20 pts)")
            total_score += 0.20
        elif total_colored_runs > 0 and new_color_count > old_color_count:
            partial = 0.10
            print(f"PARTIAL: Component 3 - {new_color_count} new vs {old_color_count} old color runs ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 - old_colors={old_color_count}, new_colors={new_color_count}, total={total_colored_runs}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Text content preserved (0.20 points)
    # All original text content must still be present after template application
    try:
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    all_text += shape.text_frame.text + "\n"

        found_count = 0
        missing = []
        for expected in EXPECTED_TEXTS:
            if expected in all_text:
                found_count += 1
            else:
                missing.append(expected)

        # Content preserved AND font changed (compound: content must be there with new styling)
        # The font change is already scored in Component 2, so here we just verify content
        # BUT to ensure this component FAILS on initial only if combined with a template change indicator,
        # we also require that the background is solid (i.e., template was applied) for credit
        has_solid_bg = any(slide.background.fill.type == 1 for slide in prs.slides)

        if found_count == len(EXPECTED_TEXTS) and has_solid_bg:
            print(f"PASS: Component 4 - All {len(EXPECTED_TEXTS)} text markers found with template applied (0.20 pts)")
            total_score += 0.20
        elif found_count == len(EXPECTED_TEXTS) and not has_solid_bg:
            # Content is there but no template applied - this is the initial state, no credit
            print(f"FAIL: Component 4 - Text preserved but template not applied (no solid bg)")
        elif found_count >= len(EXPECTED_TEXTS) * 0.8:
            partial = 0.10
            print(f"PARTIAL: Component 4 - {found_count}/{len(EXPECTED_TEXTS)} text markers found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 - Only {found_count}/{len(EXPECTED_TEXTS)} text markers found. Missing: {missing[:5]}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    persist_app_state("libreoffice_impress")
    verify_task(FILE_PATH)
