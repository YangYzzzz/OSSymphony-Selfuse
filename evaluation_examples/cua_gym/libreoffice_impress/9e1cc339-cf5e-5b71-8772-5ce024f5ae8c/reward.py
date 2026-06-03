"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 232, the title is getting lost against the background. I need to switch its font color to Purple 4 (#551A8B) and add a 0.5 pt outline so it stands out. How can I do that in LibreOffice Impress?
Generated: 2025-09-10 19:56:55
Status: success
Model: azure-o3
Total Steps: 5
"""

import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

def verify_impress_task(file_path: str) -> float:
    """Verify that on slide 232 the title font color is Purple 4 (#551A8B)
    and a 0.5 pt outline has been applied.  Returns a progressive score
    between 0.0 and 1.0 (float). Prints detailed diagnostics and the final
    reward in the required format.
    """

    # Weighting for progressive scoring
    FONT_COLOR_POINTS   = 0.5   # points for correct font colour
    OUTLINE_WIDTH_POINTS = 0.5  # points for correct outline width
    MAX_SCORE = 1.0

    target_rgb   = RGBColor(0x55, 0x1A, 0x8B)  # Purple 4 – #551A8B
    target_width = 6350                         # 0.5 pt in EMUs (1 pt = 12700 EMU)
    width_tol    = 300                          # ±0.024 pt tolerance

    print(f"Verifying presentation: {file_path}")

    # ------------------------------------------------------------------
    # 1. Safety checks – file exists & loads
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task")
        print("REWARD: 0.0")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to load presentation: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Ensure slide 232 exists (index 231)
    # ------------------------------------------------------------------
    slide_idx = 231  # zero-based index for the 232nd slide
    if len(prs.slides) <= slide_idx:
        print(f"✗ Slide 232 not present (only {len(prs.slides)} slides)")
        print("REWARD: 0.0")
        return 0.0
    slide = prs.slides[slide_idx]

    # ------------------------------------------------------------------
    # 3. Locate the title placeholder (TITLE or CENTER_TITLE)
    # ------------------------------------------------------------------
    title_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            title_shape = shape
            break

    if title_shape is None or not title_shape.has_text_frame:
        print("✗ Title placeholder with text frame not found on slide 232")
        print("REWARD: 0.0")
        return 0.0
    print("✓ Found title placeholder on slide 232")

    # Gather all runs in the title text
    runs = [run for para in title_shape.text_frame.paragraphs for run in para.runs]
    if not runs:
        print("✗ Title textbox contains no text runs – nothing to verify")
        print("REWARD: 0.0")
        return 0.0

    score = 0.0  # progressive score accumulator

    # ------------------------------------------------------------------
    # 4. Verify font colour on every run (Purple 4)
    # ------------------------------------------------------------------
    correct_runs = 0
    for run in runs:
        try:
            if run.font.color.type == 1 and run.font.color.rgb == target_rgb:  # 1 == RGB colour
                correct_runs += 1
        except Exception:
            # If colour type is not RGB or inaccessible, treat as incorrect
            pass
    if correct_runs == len(runs):
        print(f"✓ All {correct_runs}/{len(runs)} text runs have the correct font colour #551A8B")
        score += FONT_COLOR_POINTS
    elif correct_runs > 0:
        partial = FONT_COLOR_POINTS * (correct_runs / len(runs))
        print(f"⚠️  {correct_runs}/{len(runs)} runs have the correct colour (partial credit {partial:.2f})")
        score += partial
    else:
        print("✗ Font colour #551A8B not applied to any runs")

    # ------------------------------------------------------------------
    # 5. Verify 0.5 pt outline on the title text
    # ------------------------------------------------------------------
    outline_ok = False

    # 5a. Shape-level outline check (Title textbox as a whole)
    shape_line_width = title_shape.line.width if title_shape.line.width is not None else 0
    if shape_line_width and abs(shape_line_width - target_width) <= width_tol:
        print(f"✓ Title shape outline width ≈0.5 pt ({shape_line_width} EMU)")
        outline_ok = True

    # 5b. Run-level text outline (individual text runs)
    if not outline_ok:
        from pptx.oxml import OxmlElement
        for run in runs:
            rPr = run._r.rPr
            ln  = rPr.find(qn('a:ln')) if rPr is not None else None
            if ln is not None and 'w' in ln.attrib:
                try:
                    w = int(ln.get('w'))
                    if abs(w - target_width) <= width_tol:
                        print(f"✓ Run-level text outline width ≈0.5 pt ({w} EMU)")
                        outline_ok = True
                        break
                except ValueError:
                    # Non-integer width – ignore
                    pass

    if outline_ok:
        score += OUTLINE_WIDTH_POINTS
    else:
        print("✗ 0.5 pt outline not detected on title text")

    # ------------------------------------------------------------------
    # 6. Finalise score and report
    # ------------------------------------------------------------------
    final_score = min(score, MAX_SCORE)
    print(f"Final score: {final_score}/{MAX_SCORE}")
    print(f"REWARD: {final_score}")
    return final_score


# ----------------------------------------------------------------------
# Execute verification when run as a script
# ----------------------------------------------------------------------
if __name__ == "__main__":
    test_path = "/home/user/on_slide_232_the_title_is_getting_lost_against_the_background_i_need_to_switch_its_font_color_to_pur_golden.pptx"
    verify_impress_task(test_path)

