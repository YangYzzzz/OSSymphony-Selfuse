"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert 'Page X of Y' centered in the footer.
Generated: 2025-10-17 15:15:28
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
import re
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def verify_page_x_of_y_footer(file_path: str) -> float:
    """Verify that every slide in the presentation has
    a centred footer containing the text 'Page X of Y'.

    Scoring (progressive):
        • 0.7 points – proportion of slides that contain the correct
          "Page X of Y" text (case-insensitive, exact numbers).
        • 0.3 points – proportion of slides where that footer text is
          centre-aligned.

    The final score is capped at 1.0 and rounded to two decimals.
    """

    # ---------- 1. Basic file checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    if total_slides == 0:
        print("✗ Presentation has no slides")
        print("REWARD: 0.0")
        return 0.0

    print(f"Total slides detected: {total_slides}\n")

    # ---------- 2. Iterate over slides and verify footer ----------
    correct_presence = 0  # slides with correct text
    correct_alignment = 0  # slides with text centred

    for slide_idx, slide in enumerate(prs.slides, start=1):
        expected_pattern = re.compile(rf"^Page\s+{slide_idx}\s+of\s+{total_slides}$", re.IGNORECASE)
        matching_shape = None

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if expected_pattern.match(text):
                matching_shape = shape
                break

        if matching_shape:
            correct_presence += 1
            # Check alignment of first paragraph (best indicator for shape)
            para_alignment = matching_shape.text_frame.paragraphs[0].alignment
            alignment_name = {
                PP_ALIGN.LEFT: "LEFT",
                PP_ALIGN.CENTER: "CENTER",
                PP_ALIGN.RIGHT: "RIGHT",
                PP_ALIGN.JUSTIFY: "JUSTIFY",
            }.get(para_alignment, str(para_alignment))

            print(f"✓ Slide {slide_idx}: footer text found (alignment={alignment_name})")

            if para_alignment == PP_ALIGN.CENTER:
                correct_alignment += 1
            else:
                print("  ✗ Footer text not centre-aligned")
        else:
            print(f"✗ Slide {slide_idx}: expected 'Page {slide_idx} of {total_slides}' not found")

    # ---------- 3. Progressive scoring ----------
    presence_ratio = correct_presence / total_slides
    alignment_ratio = correct_alignment / total_slides

    presence_score = presence_ratio * 0.7  # 70 % weight
    alignment_score = alignment_ratio * 0.3  # 30 % weight

    total_score = round(min(presence_score + alignment_score, 1.0), 2)

    print("\n---------- SCORE BREAKDOWN ----------")
    print(f"Presence ratio  : {presence_ratio:.2f} ⇒ {presence_score:.2f} points")
    print(f"Alignment ratio : {alignment_ratio:.2f} ⇒ {alignment_score:.2f} points")
    print("------------------------------------")
    print(f"REWARD: {total_score}")

    return total_score


if __name__ == "__main__":
    # Path provided by task context
    FILE_PATH = "/home/user/insert_page_x_of_y_centered_in_the_footer.pptx"
    verify_page_x_of_y_footer(FILE_PATH)

