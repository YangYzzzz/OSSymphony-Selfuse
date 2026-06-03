"""
Initial Setup: Add company logo to master slide, insert footer and fixed date on slides 2-12
Task ID: impress_gf4_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_PATH = f'{WORKDIR}/Desktop/company_logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_company_logo():
    """Create a realistic company logo PNG on the Desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    img = Image.new('RGBA', (200, 200), (255, 255, 255, 0))
    draw = ImageDraw.Draw(img)
    # Blue circle background
    draw.ellipse([10, 10, 190, 190], fill=(0, 102, 204, 255))
    # White "AB" letters as logo
    draw.text((55, 60), "AB", fill=(255, 255, 255, 255))
    # Orange accent bar
    draw.rectangle([40, 140, 160, 155], fill=(255, 153, 0, 255))
    img.save(LOGO_PATH)
    print(f'Company logo created: {LOGO_PATH}')


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Internal Briefing"
    slide1.placeholders[1].text = "Q4 2024 Strategic Review\nAlphabridge Consulting"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Financial Performance Overview"
    items = [
        "Regional Market Analysis",
        "Client Portfolio Updates",
        "Technology Infrastructure Roadmap",
        "Talent Acquisition & Retention",
        "Risk Management Assessment",
        "2025 Strategic Priorities",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Financial Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Financial Performance"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Revenue: $12.4M (up 18% YoY)"
    for line in [
        "Operating Margin: 24.3% (target: 22%)",
        "New Client Revenue: $3.1M",
        "Client Retention Rate: 94.7%",
        "EBITDA: $3.8M",
    ]:
        p = body3.add_paragraph()
        p.text = line

    # --- Slide 4: Regional Analysis ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Market Analysis"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North America: Strong growth in fintech vertical (+32%)"
    for line in [
        "EMEA: Stable performance, new Frankfurt office opened",
        "APAC: Singapore hub expanded, revenue up 41%",
        "LATAM: Pilot program in Sao Paulo showing promise",
    ]:
        p = body4.add_paragraph()
        p.text = line

    # --- Slide 5: Client Portfolio ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Client Portfolio Updates"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Top 10 clients represent 58% of revenue"
    for line in [
        "3 new enterprise accounts signed in Q4",
        "Average contract value increased to $420K",
        "NPS score: 72 (industry avg: 54)",
    ]:
        p = body5.add_paragraph()
        p.text = line

    # --- Slide 6: Technology Roadmap ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Technology Infrastructure"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Cloud migration 85% complete (target: Q1 2025)"
    for line in [
        "AI-driven analytics platform launched in beta",
        "Security audit: zero critical findings",
        "DevOps pipeline reduced deployment time by 60%",
    ]:
        p = body6.add_paragraph()
        p.text = line

    # --- Slide 7: Talent ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Talent Acquisition & Retention"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Headcount: 187 (+23 in Q4)"
    for line in [
        "Voluntary attrition: 8.2% (down from 11.4%)",
        "Senior hires: VP of Engineering, Director of Data Science",
        "Employee satisfaction score: 4.3/5.0",
    ]:
        p = body7.add_paragraph()
        p.text = line

    # --- Slide 8: Risk Management ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Risk Management Assessment"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Regulatory compliance: fully compliant across all jurisdictions"
    for line in [
        "Cybersecurity insurance coverage increased to $10M",
        "Business continuity plan tested successfully",
        "Vendor concentration risk: moderate (action plan in place)",
    ]:
        p = body8.add_paragraph()
        p.text = line

    # --- Slide 9: Strategic Priorities ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "2025 Strategic Priorities"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "1. Expand APAC operations to Tokyo and Sydney"
    for line in [
        "2. Launch managed services offering",
        "3. Achieve ISO 27001 certification",
        "4. Grow revenue to $16M",
        "5. Develop partner ecosystem program",
    ]:
        p = body9.add_paragraph()
        p.text = line

    # --- Slide 10: Budget Allocation ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "2025 Budget Allocation"
    body10 = slide10.placeholders[1].text_frame
    body10.text = "Engineering & Product: 35% ($5.6M)"
    for line in [
        "Sales & Marketing: 28% ($4.5M)",
        "Operations: 20% ($3.2M)",
        "R&D Innovation Fund: 12% ($1.9M)",
        "Corporate Overhead: 5% ($0.8M)",
    ]:
        p = body10.add_paragraph()
        p.text = line

    # --- Slide 11: Timeline ---
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Key Milestones Timeline"
    body11 = slide11.placeholders[1].text_frame
    body11.text = "Q1 2025: Tokyo office launch, ISO audit Phase 1"
    for line in [
        "Q2 2025: Managed services beta, Partner program kickoff",
        "Q3 2025: Sydney office setup, AI platform GA release",
        "Q4 2025: Annual review, ISO certification expected",
    ]:
        p = body11.add_paragraph()
        p.text = line

    # --- Slide 12: Thank You / Q&A ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[0])
    slide12.shapes.title.text = "Thank You"
    slide12.placeholders[1].text = "Questions & Discussion\nContact: strategy@alphabridge.com"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create the logo first, then the presentation
create_company_logo()
create_initial()
