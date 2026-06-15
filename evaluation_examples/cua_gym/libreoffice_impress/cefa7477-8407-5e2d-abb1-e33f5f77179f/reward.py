"""
Reward Script: Master slide logo, footer, and date field verification
Task ID: impress_gf4_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.35): Company logo image present on master slide in top-right corner
  Component 2 (0.35): Footer 'Confidential - Internal Use Only' on slides 2-12 but NOT slide 1
  Component 3 (0.30): Date '2024-01-15' on slides 2-12 but NOT slide 1
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_007'


def persist_app_state(domain: str):
    """Best-effort save of any open LibreOffice document."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 2:
        print(f"CRITICAL: Expected at least 12 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ======================================================================
    # Component 1: Company logo on master slide in top-right corner (0.35)
    # This FAILS on initial (no image on master) and PASSES on golden.
    # ======================================================================
    try:
        logo_found = False
        logo_top_right = False
        logo_blob_match = False

        master = prs.slide_masters[0]
        for shape in master.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                logo_found = True
                # Check position: top-right means left > midpoint and top < midpoint
                mid_x = slide_width / 2
                mid_y = slide_height / 2
                if shape.left > mid_x and shape.top < mid_y:
                    logo_top_right = True

                # Check if blob matches company_logo.png
                logo_path = os.path.join(WORKDIR, 'Desktop', 'company_logo.png')
                if os.path.exists(logo_path):
                    with open(logo_path, 'rb') as f:
                        expected_blob = f.read()
                    if shape.image.blob == expected_blob:
                        logo_blob_match = True

                break  # Only check first picture on master

        if logo_found and logo_top_right and logo_blob_match:
            print(f"PASS: Component 1 — Logo on master, top-right, blob matches (0.35 pts)")
            total_score += 0.35
        elif logo_found and logo_top_right:
            # Logo exists in right place but different image — partial
            print(f"PASS (partial): Component 1 — Logo on master in top-right but blob mismatch (0.25 pts)")
            total_score += 0.25
        elif logo_found:
            # Logo exists but wrong position
            print(f"PASS (partial): Component 1 — Logo on master but not in top-right (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — No image found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ======================================================================
    # Component 2: Footer text on slides 2-12 but NOT on slide 1 (0.35)
    # This FAILS on initial (no footer textboxes) and PASSES on golden.
    # ======================================================================
    try:
        expected_footer = 'Confidential \u2013 Internal Use Only'
        footer_on_slides_2_plus = 0
        footer_on_slide_1 = False

        # Check slide 1 (index 0)
        for shape in prs.slides[0].shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if 'Confidential' in txt and 'Internal Use Only' in txt:
                    footer_on_slide_1 = True
                    break

        # Check slides 2-12 (indices 1 through num_slides-1)
        expected_footer_slides = num_slides - 1  # slides 2 through N
        for idx in range(1, num_slides):
            slide = prs.slides[idx]
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if 'Confidential' in txt and 'Internal Use Only' in txt:
                        footer_on_slides_2_plus += 1
                        break

        # Score: footer on all slides 2+ AND not on slide 1
        if footer_on_slides_2_plus == expected_footer_slides and not footer_on_slide_1:
            print(f"PASS: Component 2 — Footer on {footer_on_slides_2_plus}/{expected_footer_slides} slides (2+), absent from slide 1 (0.35 pts)")
            total_score += 0.35
        elif footer_on_slides_2_plus > 0 and not footer_on_slide_1:
            # Partial: some slides have footer
            partial = 0.35 * (footer_on_slides_2_plus / expected_footer_slides)
            print(f"PASS (partial): Component 2 — Footer on {footer_on_slides_2_plus}/{expected_footer_slides} slides ({partial:.2f} pts)")
            total_score += partial
        elif footer_on_slides_2_plus > 0 and footer_on_slide_1:
            # Footer on slide 1 too — less credit
            partial = 0.35 * (footer_on_slides_2_plus / expected_footer_slides) * 0.5
            print(f"PARTIAL: Component 2 — Footer present but also on slide 1 ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Footer found on {footer_on_slides_2_plus}/{expected_footer_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ======================================================================
    # Component 3: Date '2024-01-15' on slides 2-12 but NOT on slide 1 (0.30)
    # This FAILS on initial (no date textboxes) and PASSES on golden.
    # ======================================================================
    try:
        expected_date = '2024-01-15'
        date_on_slides_2_plus = 0
        date_on_slide_1 = False

        # Check slide 1 (index 0)
        for shape in prs.slides[0].shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                txt = shape.text_frame.text.strip()
                if expected_date in txt:
                    date_on_slide_1 = True
                    break

        # Check slides 2-12
        expected_date_slides = num_slides - 1
        for idx in range(1, num_slides):
            slide = prs.slides[idx]
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if expected_date in txt:
                        date_on_slides_2_plus += 1
                        break

        if date_on_slides_2_plus == expected_date_slides and not date_on_slide_1:
            print(f"PASS: Component 3 — Date '{expected_date}' on {date_on_slides_2_plus}/{expected_date_slides} slides (2+), absent from slide 1 (0.30 pts)")
            total_score += 0.30
        elif date_on_slides_2_plus > 0 and not date_on_slide_1:
            partial = 0.30 * (date_on_slides_2_plus / expected_date_slides)
            print(f"PASS (partial): Component 3 — Date on {date_on_slides_2_plus}/{expected_date_slides} slides ({partial:.2f} pts)")
            total_score += partial
        elif date_on_slides_2_plus > 0 and date_on_slide_1:
            partial = 0.30 * (date_on_slides_2_plus / expected_date_slides) * 0.5
            print(f"PARTIAL: Component 3 — Date present but also on slide 1 ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — Date '{expected_date}' found on {date_on_slides_2_plus}/{expected_date_slides} slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
