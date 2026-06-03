"""
Reward Script: Add a new blank slide at the end of the presentation
Task ID: impstruct_001
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Presentation has exactly 5 slides
  Component 2 (0.3): 5th slide uses "Blank" layout
  Component 3 (0.3): Original 4 slides preserved with correct titles
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impstruct_001'

# Expected titles for slides 1-4 (from initial state)
EXPECTED_SLIDE_TITLES = [
    'Q3 Results',
    'Revenue Overview',
    'Expenses',
    'Outlook',
]


def get_slide_title(slide):
    """Extract the title text from a slide (first text shape with content)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return ''


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 5 slides (0.4 points)
    # Initial has 4 slides; golden should have 5
    try:
        if num_slides == 5:
            print(f"PASS: Component 1 — Slide count is 5 (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Expected 5 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: 5th slide uses "Blank" layout (0.3 points)
    # Only check if there are at least 5 slides
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            layout_name = slide5.slide_layout.name if slide5.slide_layout else 'None'
            if layout_name == 'Blank':
                print(f"PASS: Component 2 — Slide 5 layout is 'Blank' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — Expected 'Blank' layout for slide 5, found '{layout_name}'")
        else:
            print(f"FAIL: Component 2 — Not enough slides to check (need >= 5, have {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Original 4 slides preserved with correct titles (0.3 points)
    # Verify that slides 1-4 still have their original title text
    try:
        if num_slides >= 5:
            mismatched = [
                i for i in range(4)
                if get_slide_title(prs.slides[i]) != EXPECTED_SLIDE_TITLES[i]
            ]
            if len(mismatched) == 0:
                print(f"PASS: Component 3 — All 4 original slide titles preserved (0.3 pts)")
                total_score += 0.3
            else:
                for i in mismatched:
                    actual_title = get_slide_title(prs.slides[i])
                    print(f"FAIL: Component 3 — Slide {i+1} title mismatch: expected '{EXPECTED_SLIDE_TITLES[i]}', found '{actual_title}'")
        else:
            print(f"FAIL: Component 3 — Not enough slides to verify preservation (need >= 5, have {num_slides})")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
