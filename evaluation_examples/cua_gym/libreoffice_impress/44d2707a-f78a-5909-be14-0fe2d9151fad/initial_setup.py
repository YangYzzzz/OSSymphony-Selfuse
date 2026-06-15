"""
Initial Setup: Add a new blank slide at the end of the presentation.
Task ID: impstruct_001
Domain: libreoffice_impress

Creates a 4-slide quarterly report presentation and opens it in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impstruct_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Q3 Results (Title Slide) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q3 Results"
    slide1.placeholders[1].text = "Quarterly Business Review\nFiscal Year 2025"

    # --- Slide 2: Revenue Overview (Title + Content) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Revenue Overview"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total Revenue: $4.82M (+12% YoY)"
    p = body2.add_paragraph()
    p.text = "Product Sales: $3.15M"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Service Revenue: $1.67M"
    p.level = 1
    p = body2.add_paragraph()
    p.text = ""
    p = body2.add_paragraph()
    p.text = "Key Highlights:"
    p = body2.add_paragraph()
    p.text = "Enterprise segment grew 18% driven by new contracts"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "APAC region exceeded targets by 7%"
    p.level = 1
    p = body2.add_paragraph()
    p.text = "Customer retention rate: 94.2%"
    p.level = 1

    # --- Slide 3: Expenses (Title + Content) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Expenses"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total Operating Expenses: $3.41M"
    p = body3.add_paragraph()
    p.text = "Personnel: $1.95M (57%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Infrastructure: $0.68M (20%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "Marketing: $0.52M (15%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = "R&D: $0.26M (8%)"
    p.level = 1
    p = body3.add_paragraph()
    p.text = ""
    p = body3.add_paragraph()
    p.text = "Operating Margin: 29.3% (up from 25.1% in Q2)"

    # --- Slide 4: Outlook (Title + Content) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Outlook"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Q4 Projections"
    p = body4.add_paragraph()
    p.text = "Revenue target: $5.20M"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Expected headcount increase: 12 FTEs"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "New product launch: CloudSync Pro (October 2025)"
    p.level = 1
    p = body4.add_paragraph()
    p.text = ""
    p = body4.add_paragraph()
    p.text = "Strategic Priorities:"
    p = body4.add_paragraph()
    p.text = "Expand enterprise partnerships in EMEA"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Achieve SOC 2 Type II certification by December"
    p.level = 1
    p = body4.add_paragraph()
    p.text = "Reduce customer onboarding time by 30%"
    p.level = 1

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
