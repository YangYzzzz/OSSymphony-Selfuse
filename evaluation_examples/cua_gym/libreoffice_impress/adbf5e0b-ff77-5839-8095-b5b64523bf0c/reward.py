"""
Reward Script: Change content alignment on slides 1, 3, and 5 to center-aligned.
Task ID: osworld_impress_per_slide_alignment_005
Domain: libreoffice_impress
Scoring:
  Component 1: Each of content placeholders on slides 1, 3, 5 is CENTER-aligned (0.2 pts each = 0.6 total)
               FAILS on initial (all LEFT), PASSES on golden (center slides changed)
  Component 2: Compound check — slides 1, 3, 5 are ALL CENTER AND slides 2, 4, 6 are ALL LEFT (0.4 pts)
               FAILS on initial (slides 1,3,5 still LEFT), PASSES on golden (correct combination)
  Total: 1.0
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_005'

# Slides that should be CENTER-aligned (1-indexed)
CENTER_SLIDES = [1, 3, 5]
# Slides that should remain LEFT-aligned (1-indexed)
LEFT_SLIDES = [2, 4, 6]


def get_content_placeholder(slide):
    """Return the content/body placeholder from a slide, or None if not found."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Content Placeholder'):
            return shape
    return None


def check_paragraphs_alignment(shape, expected_alignment):
    """
    Check that ALL non-empty paragraphs in the shape's text_frame have the expected alignment.
    Returns (all_match: bool, details: str).
    Note: PP_ALIGN.LEFT == 1, PP_ALIGN.CENTER == 2.
    None alignment value defaults to LEFT.
    """
    tf = shape.text_frame
    mismatches = []
    checked = 0
    for para_idx, para in enumerate(tf.paragraphs):
        text_preview = para.text[:30] if para.text else ''
        if not para.text.strip():
            continue  # skip empty paragraphs
        checked += 1
        actual = para.alignment
        # Normalize: None == LEFT
        if actual is None:
            actual = PP_ALIGN.LEFT
        if actual != expected_alignment:
            mismatches.append(
                f"Para {para_idx} ({text_preview!r}): expected {expected_alignment}, got {actual}"
            )
    if checked == 0:
        return False, "No non-empty paragraphs found"
    if mismatches:
        return False, "; ".join(mismatches)
    return True, f"All {checked} paragraphs correct"


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Presentation has {len(prs.slides)} slides")

    # Component 1: Content placeholders on slides 1, 3, 5 are CENTER-aligned (0.2 pts each = 0.6 total)
    # Each sub-check FAILS on initial_env (all LEFT) and PASSES on golden_env (changed to CENTER)
    center_slide_results = []
    try:
        for slide_num in CENTER_SLIDES:
            slide = prs.slides[slide_num - 1]  # 0-indexed
            content_shape = get_content_placeholder(slide)
            if content_shape is None:
                center_slide_results.append((slide_num, False, "No content placeholder found"))
                continue
            ok, details = check_paragraphs_alignment(content_shape, PP_ALIGN.CENTER)
            center_slide_results.append((slide_num, ok, details))

        for slide_num, ok, details in center_slide_results:
            if ok:
                print(f"PASS: Slide {slide_num} content is CENTER-aligned — {details} (+0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Slide {slide_num} content not CENTER-aligned — {details}")
        print(f"Component 1 subtotal: {sum(0.2 for _, ok, _ in center_slide_results if ok):.1f}/0.6")
    except Exception as e:
        print(f"ERROR: Component 1 (center slides check) — {e}")

    # Component 2: Compound check — slides 1, 3, 5 are ALL CENTER AND slides 2, 4, 6 are ALL LEFT (0.4 pts)
    # This compound check FAILS on initial_env because slides 1, 3, 5 are still LEFT (not CENTER),
    # and PASSES on golden_env because both conditions are satisfied simultaneously.
    # This rewards the "correct selective alignment" — changed the right slides AND preserved the others.
    try:
        # Check that center slides are ALL center (re-use results from Component 1)
        all_center_ok = all(ok for _, ok, _ in center_slide_results)

        # Check that left slides are ALL left
        left_slide_results = []
        for slide_num in LEFT_SLIDES:
            slide = prs.slides[slide_num - 1]
            content_shape = get_content_placeholder(slide)
            if content_shape is None:
                left_slide_results.append((slide_num, False, "No content placeholder found"))
                continue
            ok, details = check_paragraphs_alignment(content_shape, PP_ALIGN.LEFT)
            left_slide_results.append((slide_num, ok, details))

        all_left_ok = all(ok for _, ok, _ in left_slide_results)

        # Only award 0.4 pts if BOTH conditions are satisfied simultaneously
        if all_center_ok and all_left_ok:
            print(f"PASS: Component 2 — Slides 1,3,5 all CENTER AND slides 2,4,6 all LEFT (+0.4 pts)")
            total_score += 0.4
        else:
            if not all_center_ok:
                print(f"FAIL: Component 2 — Not all target slides (1,3,5) are CENTER-aligned")
            if not all_left_ok:
                untouched_failures = [(s, d) for s, ok, d in left_slide_results if not ok]
                for slide_num, details in untouched_failures:
                    print(f"FAIL: Component 2 — Slide {slide_num} should remain LEFT — {details}")
        print(f"Component 2 subtotal: {0.4 if all_center_ok and all_left_ok else 0.0:.1f}/0.4")
    except Exception as e:
        print(f"ERROR: Component 2 (compound alignment check) — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
