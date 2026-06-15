"""
Initial Setup: 6-slide training materials deck with all content textboxes left-aligned.
Task ID: osworld_impress_per_slide_alignment_005
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_005'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_paragraph_left(para):
    """Explicitly set a paragraph to left alignment."""
    para.alignment = PP_ALIGN.LEFT


def add_content_slide(prs, title_text, body_lines):
    """Add a slide with a title and left-aligned body content."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text

    # Set body content with explicit LEFT alignment
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(18)

    return slide


def create_initial():
    prs = Presentation()
    # Use standard widescreen dimensions (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Course Introduction ---
    add_content_slide(
        prs,
        "Welcome to New Employee Onboarding",
        [
            "This training program is designed to help you get started at Meridian Solutions.",
            "You will learn about our company values, policies, and key tools.",
            "Duration: 3 days of instructor-led sessions plus self-paced modules.",
            "Please review all materials before the end of your first week.",
            "Contact HR at hr@meridiansolutions.com with any questions.",
        ]
    )

    # --- Slide 2: Company Overview ---
    add_content_slide(
        prs,
        "About Meridian Solutions",
        [
            "Founded in 2008, Meridian Solutions serves clients across 14 countries.",
            "Core business areas: Cloud Infrastructure, Data Analytics, Cybersecurity.",
            "Over 3,200 employees worldwide with headquarters in Singapore.",
            "Revenue: $850M (FY2024) — 18% YoY growth.",
            "Mission: Empowering organizations through intelligent technology.",
        ]
    )

    # --- Slide 3: Code of Conduct ---
    add_content_slide(
        prs,
        "Code of Conduct & Workplace Policies",
        [
            "All employees must adhere to the Meridian Code of Conduct.",
            "Key principles: Integrity, Respect, Accountability, Innovation.",
            "Zero tolerance policy for harassment, discrimination, and fraud.",
            "Confidentiality agreements must be signed before system access is granted.",
            "Annual compliance training is mandatory for all staff.",
        ]
    )

    # --- Slide 4: IT & Security Setup ---
    add_content_slide(
        prs,
        "IT Systems & Security Requirements",
        [
            "Submit your IT equipment request form via the Employee Portal within 24 hours.",
            "Multi-factor authentication (MFA) is required on all corporate accounts.",
            "Use only approved VPN software when accessing internal systems remotely.",
            "Report suspicious emails to security@meridiansolutions.com immediately.",
            "Password policy: minimum 12 characters, rotate every 90 days.",
        ]
    )

    # --- Slide 5: Benefits & Compensation ---
    add_content_slide(
        prs,
        "Employee Benefits & Compensation",
        [
            "Health insurance: comprehensive coverage starting from Day 1 of employment.",
            "Annual leave: 18 days per year (prorated in first year).",
            "Performance bonuses reviewed bi-annually in June and December.",
            "Professional development allowance: $1,500 per year per employee.",
            "Employee stock purchase plan (ESPP) available after 6 months tenure.",
        ]
    )

    # --- Slide 6: Next Steps & Resources ---
    add_content_slide(
        prs,
        "Next Steps & Key Resources",
        [
            "Complete your profile on the Meridian People portal by end of Day 1.",
            "Schedule your 1-on-1 meeting with your manager within the first week.",
            "Enroll in required compliance courses via the Learning Management System.",
            "Join the #new-employees Slack channel for peer support and announcements.",
            "Your 30/60/90 day plan will be reviewed with your manager at each milestone.",
        ]
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
