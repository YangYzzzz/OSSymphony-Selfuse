"""
Reward Script: Insert SmartArt-style hierarchy diagram on slide 3
Task ID: impress_stu_042
Domain: libreoffice_impress
Scoring:
  Component 1 (0.2): Rounded rectangles present on slide 3
  Component 2 (0.4): All 7 hierarchy node texts present
  Component 3 (0.2): Connecting lines/connectors present
  Component 4 (0.2): Hierarchical vertical layout (3 levels)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_042'

# Required hierarchy node texts (case-insensitive matching)
REQUIRED_TEXTS = {
    'machine learning',
    'supervised',
    'unsupervised',
    'classification',
    'regression',
    'clustering',
    'dimensionality reduction',
}

# Expected hierarchy levels (by text) for layout verification
# Level 1 (top): Machine Learning
# Level 2 (middle): Supervised, Unsupervised
# Level 3 (bottom): Classification, Regression, Clustering, Dimensionality Reduction
LEVEL_1_TEXTS = {'machine learning'}
LEVEL_2_TEXTS = {'supervised', 'unsupervised'}
LEVEL_3_TEXTS = {'classification', 'regression', 'clustering', 'dimensionality reduction'}


def get_all_shapes_recursive(slide):
    """Get all shapes including those inside groups."""
    results = []
    for shape in slide.shapes:
        results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.append(sub)
    return results


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # Slide 3 (0-indexed)
    all_shapes = get_all_shapes_recursive(slide)

    # Collect rounded rectangles and their properties
    rounded_rects = []
    for shape in all_shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                    text = shape.text.replace(chr(11), ' ').strip()
                    rounded_rects.append({
                        'text': text,
                        'text_lower': text.lower(),
                        'top': shape.top,
                        'left': shape.left,
                    })
        except Exception:
            pass

    # Also check for any auto shapes that might contain the right text
    # (in case agent uses different shape types)
    auto_shapes_with_text = []
    for shape in all_shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(shape, 'text'):
                text = shape.text.replace(chr(11), ' ').strip()
                if text:
                    auto_shapes_with_text.append({
                        'text': text,
                        'text_lower': text.lower(),
                        'top': shape.top,
                        'left': shape.left,
                    })
        except Exception:
            pass

    # Count connectors/lines
    connector_count = 0
    for shape in all_shapes:
        try:
            stype = shape.shape_type
            if stype == 9:  # LINE
                connector_count += 1
            elif stype == MSO_SHAPE_TYPE.FREEFORM:
                # Some connectors may be freeforms
                connector_count += 1
        except Exception:
            pass

    # Component 1: Rounded rectangles present on slide 3 (0.2 points)
    # Task specifically asks for "rounded rectangles"
    try:
        num_rects = len(rounded_rects)
        if num_rects >= 7:
            print(f"PASS: Component 1 -- {num_rects} rounded rectangles found on slide 3 (0.2 pts)")
            total_score += 0.2
        elif num_rects >= 4:
            partial = 0.1
            print(f"PARTIAL: Component 1 -- {num_rects}/7 rounded rectangles found (0.1 pts)")
            total_score += partial
        else:
            # Fall back: check if there are auto shapes with hierarchy text
            if len(auto_shapes_with_text) >= 7:
                print(f"PARTIAL: Component 1 -- {len(auto_shapes_with_text)} auto shapes (not all rounded rects) found (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 1 -- Only {num_rects} rounded rectangles found, expected >= 7")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All 7 hierarchy node texts present (0.4 points)
    # Primary: check rounded rects; fallback: check all auto shapes
    try:
        # Prefer rounded_rects; fall back to auto_shapes_with_text only if more matches
        shapes_to_check = rounded_rects if rounded_rects else auto_shapes_with_text
        found_texts = set()
        for rect in shapes_to_check:
            # Normalize: lowercase, collapse whitespace
            t = ' '.join(rect['text_lower'].split())
            for required in REQUIRED_TEXTS:
                # Exact match only (after normalization) to avoid false positives
                # e.g., "Types of Machine Learning" should NOT match "machine learning"
                if t == required:
                    found_texts.add(required)

        num_found = len(found_texts)
        if num_found == 7:
            print(f"PASS: Component 2 -- All 7 hierarchy texts found: {found_texts} (0.4 pts)")
            total_score += 0.4
        elif num_found >= 1:
            partial = round(0.4 * num_found / 7, 2)
            missing = REQUIRED_TEXTS - found_texts
            print(f"PARTIAL: Component 2 -- {num_found}/7 texts found, missing: {missing} ({partial} pts)")
            if partial > 0:
                total_score += partial
        else:
            print(f"FAIL: Component 2 -- No hierarchy texts found in auto shapes on slide 3")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Connecting lines/connectors present (0.2 points)
    # Task asks for "connecting lines" between hierarchy nodes
    # Expected: 6 connectors (ML->Supervised, ML->Unsupervised,
    #   Supervised->Classification, Supervised->Regression,
    #   Unsupervised->Clustering, Unsupervised->DimReduction)
    try:
        if connector_count >= 6:
            print(f"PASS: Component 3 -- {connector_count} connectors/lines found (0.2 pts)")
            total_score += 0.2
        elif connector_count >= 3:
            partial = 0.1
            print(f"PARTIAL: Component 3 -- {connector_count}/6 connectors found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- Only {connector_count} connectors found, expected >= 6")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Hierarchical vertical layout (0.2 points)
    # Machine Learning should be above Supervised/Unsupervised,
    # which should be above the 4 leaf nodes
    try:
        shapes_to_check = rounded_rects if rounded_rects else auto_shapes_with_text
        level_tops = {1: [], 2: [], 3: []}
        for rect in shapes_to_check:
            t = ' '.join(rect['text_lower'].split())
            if t in LEVEL_1_TEXTS:
                level_tops[1].append(rect['top'])
            elif t in LEVEL_2_TEXTS:
                level_tops[2].append(rect['top'])
            elif t in LEVEL_3_TEXTS:
                level_tops[3].append(rect['top'])

        if level_tops[1] and level_tops[2] and level_tops[3]:
            avg_top_1 = sum(level_tops[1]) / len(level_tops[1])
            avg_top_2 = sum(level_tops[2]) / len(level_tops[2])
            avg_top_3 = sum(level_tops[3]) / len(level_tops[3])

            # In pptx, top increases downward. So level 1 < level 2 < level 3
            if avg_top_1 < avg_top_2 < avg_top_3:
                print(f"PASS: Component 4 -- Hierarchical layout verified: L1 top={avg_top_1:.0f} < L2 top={avg_top_2:.0f} < L3 top={avg_top_3:.0f} (0.2 pts)")
                total_score += 0.2
            elif avg_top_1 < avg_top_3:
                print(f"PARTIAL: Component 4 -- Partial hierarchy: L1={avg_top_1:.0f}, L2={avg_top_2:.0f}, L3={avg_top_3:.0f} (0.1 pts)")
                total_score += 0.1
            else:
                print(f"FAIL: Component 4 -- No hierarchical layout: L1={avg_top_1:.0f}, L2={avg_top_2:.0f}, L3={avg_top_3:.0f}")
        else:
            missing_levels = [k for k, v in level_tops.items() if not v]
            print(f"FAIL: Component 4 -- Cannot verify hierarchy, missing levels: {missing_levels}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
