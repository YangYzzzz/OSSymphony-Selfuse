"""
Initial Setup: ML Overview presentation with 7 slides, slide 3 empty body
Task ID: impress_stu_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Set title and add body text lines to a slide."""
    slide.shapes.title.text = title_text
    if body_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Machine Learning Overview"
    slide1.placeholders[1].text = "A Comprehensive Introduction to ML Concepts and Methods"

    # Slide 2: Introduction to ML (layout 1 = Title + Content)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Introduction to Machine Learning", [
        "Machine learning is a subset of artificial intelligence",
        "Systems learn and improve from experience without explicit programming",
        "Applications span healthcare, finance, autonomous vehicles, and NLP",
        "Global ML market projected to reach $209 billion by 2029",
        "Key enablers: big data, computational power, and algorithmic advances",
    ])

    # Slide 3: Types of ML — EMPTY body, title only (layout 6 = Title Only)
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Types of Machine Learning"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Slide 4: Supervised Learning (layout 1)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Supervised Learning", [
        "Trained on labeled datasets with known input-output pairs",
        "Classification: predicting discrete categories (spam detection, image recognition)",
        "Regression: predicting continuous values (house prices, stock forecasting)",
        "Common algorithms: Random Forest, SVM, Neural Networks, Linear Regression",
        "Requires careful feature engineering and cross-validation",
    ])

    # Slide 5: Unsupervised Learning (layout 1)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Unsupervised Learning", [
        "Works with unlabeled data to discover hidden patterns",
        "Clustering: grouping similar data points (K-Means, DBSCAN, Hierarchical)",
        "Dimensionality Reduction: simplifying data while preserving structure (PCA, t-SNE)",
        "Association Rule Learning: discovering relationships in transactional data",
        "Widely used in customer segmentation and anomaly detection",
    ])

    # Slide 6: Model Evaluation (layout 1)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Model Evaluation Metrics", [
        "Accuracy, Precision, Recall, and F1-Score for classification",
        "Mean Squared Error (MSE) and R-squared for regression",
        "Cross-validation prevents overfitting and provides robust estimates",
        "Confusion matrix visualizes true/false positives and negatives",
        "ROC-AUC curve measures discriminative ability across thresholds",
    ])

    # Slide 7: Conclusion (layout 1)
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Conclusion & Next Steps", [
        "Machine learning continues to transform industries worldwide",
        "Choosing the right algorithm depends on data type and problem context",
        "Ethical AI and interpretability are growing areas of focus",
        "Reinforcement learning and deep learning are expanding frontiers",
        "Recommended: start with structured projects on Kaggle or UCI datasets",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
