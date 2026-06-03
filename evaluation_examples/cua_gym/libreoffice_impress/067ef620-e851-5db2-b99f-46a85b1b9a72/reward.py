"""
Reward Script: Museum Exhibition Guide Presentation
Task ID: impress_wf_090
Domain: libreoffice_impress
Scoring:
  C1 (0.15) - File exists at Desktop + exactly 10 slides
  C2 (0.10) - Slide 1 title contains 'Ancient Egypt Exhibition Guide'
  C3 (0.10) - Slide 2 has rectangles (rooms) connected by line/connector shapes
  C4 (0.10) - Slide 3 has arrow shape + era block rectangles (timeline)
  C5 (0.20) - Slides 4-8 each have 3 rectangle placeholders + callout shape
  C6 (0.10) - Slides 4-8 have Appear animations (timing XML)
  C7 (0.10) - Slide 9 has quiz two-column layout with connecting lines
  C8 (0.05) - Slide 10 has gift shop items and museum info
  C9 (0.10) - Colors #C9A959 and #263238 used in the presentation
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_090'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Egypt_Exhibition.pptx')

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def get_all_text(shape):
    """Get all text from a shape's text frame."""
    if shape.has_text_frame:
        return shape.text_frame.text
    return ""


def get_shape_colors(prs):
    """Collect all RGB colors used in text runs and shape fills."""
    colors = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            # Text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        try:
                            if run.font.color.type is not None:
                                colors.add(str(run.font.color.rgb).upper())
                        except Exception:
                            pass
            # Shape fill colors
            try:
                if hasattr(shape, 'fill') and shape.fill.type is not None:
                    if shape.fill.type == 1:  # solid fill
                        colors.add(str(shape.fill.fore_color.rgb).upper())
            except Exception:
                pass
    return colors


def check_animations(pptx_path, slide_indices):
    """Check if slides have animation timing elements (Appear animations).
    slide_indices are 1-based slide numbers.
    Returns count of slides that have animations."""
    P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    count = 0
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            for slide_num in slide_indices:
                fname = f'ppt/slides/slide{slide_num}.xml'
                try:
                    with zf.open(fname) as f:
                        root = ET.parse(f).getroot()
                        # Look for entrance animation cTn elements (presetClass="entr")
                        ctns = root.findall(f'.//{{{P_NS}}}cTn')
                        entr_count = sum(1 for c in ctns if c.get('presetClass') == 'entr')
                        if entr_count > 0:
                            count += 1
                except KeyError:
                    pass
    except Exception:
        pass
    return count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # ---- Component 1: File at correct path + 10 slides (0.15 pts) ----
    try:
        if num_slides == 10:
            print(f"PASS: Component 1 — 10 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---- Component 2: Slide 1 title 'Ancient Egypt Exhibition Guide' (0.10 pts) ----
    try:
        if num_slides >= 1:
            slide1 = prs.slides[0]
            slide1_text = " ".join(get_all_text(s) for s in slide1.shapes)
            if "Ancient Egypt Exhibition Guide" in slide1_text:
                print(f"PASS: Component 2 — Slide 1 title found (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — 'Ancient Egypt Exhibition Guide' not found in slide 1 text: {slide1_text[:150]}")
        else:
            print("FAIL: Component 2 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---- Component 3: Slide 2 has rooms (rectangles) + pathways (connectors/lines) (0.10 pts) ----
    try:
        if num_slides >= 2:
            slide2 = prs.slides[1]
            rects = []
            lines = []
            for shape in slide2.shapes:
                stype = shape.shape_type
                name = shape.name
                # Rectangles (AUTO_SHAPE with Rectangle in name, not Rounded)
                if stype == MSO_SHAPE_TYPE.AUTO_SHAPE and 'Rectangle' in name and 'Rounded' not in name:
                    rects.append(shape)
                # Lines/Connectors
                if 'Connector' in name or 'Line' in name or stype == MSO_SHAPE_TYPE.LINE:
                    lines.append(shape)

            has_rooms = len(rects) >= 3
            has_pathways = len(lines) >= 2
            if has_rooms and has_pathways:
                print(f"PASS: Component 3 — Slide 2 floor plan: {len(rects)} rooms, {len(lines)} pathways (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Slide 2 floor plan: {len(rects)} rooms (need >=3), {len(lines)} pathways (need >=2)")
        else:
            print("FAIL: Component 3 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ---- Component 4: Slide 3 has arrow + era blocks (timeline) (0.10 pts) ----
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            arrow_count = sum(1 for s in slide3.shapes if 'arrow' in s.name.lower())
            era_blocks = 0
            for shape in slide3.shapes:
                name = shape.name.lower()
                # Era blocks are rectangles with dynasty/era text
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and 'rectangle' in name and 'rounded' not in name:
                    text = get_all_text(shape).lower()
                    if 'bc' in text or 'kingdom' in text or 'dynast' in text or 'period' in text:
                        era_blocks += 1

            if arrow_count >= 1 and era_blocks >= 3:
                print(f"PASS: Component 4 — Slide 3 timeline: arrow found, {era_blocks} era blocks (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — Slide 3 timeline: arrows={arrow_count}, era_blocks={era_blocks} (need >=3)")
        else:
            print("FAIL: Component 4 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # ---- Component 5: Slides 4-8 each have 3 rect placeholders + callout (0.20 pts) ----
    # 0.04 per slide
    try:
        if num_slides >= 8:
            c5_score = 0.0
            for slide_idx in range(3, 8):  # 0-indexed: slides 4-8
                slide = prs.slides[slide_idx]
                rects = []
                callout_count = 0
                for shape in slide.shapes:
                    name = shape.name
                    stype = shape.shape_type
                    if stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if 'Rectangle' in name and 'Rounded' not in name:
                            rects.append(shape)
                        if 'Rounded' in name:
                            callout_count += 1

                if len(rects) >= 3 and callout_count >= 1:
                    c5_score += 0.04
                    print(f"  Slide {slide_idx+1}: PASS — {len(rects)} rect placeholders + callout")
                else:
                    print(f"  Slide {slide_idx+1}: FAIL — {len(rects)} rects (need >=3), callouts={callout_count}")

            if c5_score > 0:
                print(f"PASS: Component 5 — room guide slides ({c5_score:.2f} of 0.20 pts)")
                total_score += c5_score
            else:
                print(f"FAIL: Component 5 — no room guide slides passed")
        else:
            print("FAIL: Component 5 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # ---- Component 6: Slides 4-8 have Appear animations (0.10 pts) ----
    try:
        anim_count = check_animations(file_path, [4, 5, 6, 7, 8])
        if anim_count >= 4:
            print(f"PASS: Component 6 — {anim_count}/5 slides have animations (0.10 pts)")
            total_score += 0.10
        elif anim_count >= 2:
            partial = 0.05
            print(f"PARTIAL: Component 6 — {anim_count}/5 slides have animations ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 6 — only {anim_count}/5 slides have animations")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # ---- Component 7: Slide 9 quiz with two columns + connecting lines (0.10 pts) ----
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            text_boxes = []
            line_shapes = []
            for shape in slide9.shapes:
                if shape.has_text_frame and shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
                    text = get_all_text(shape).strip()
                    if text:
                        text_boxes.append(text)
                if 'Connector' in shape.name or 'Line' in shape.name or shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    line_shapes.append(shape)

            # Should have numbered items (left col) + lettered items (right col) + lines
            has_numbered = sum(1 for t in text_boxes if t and t[0].isdigit()) >= 3
            has_lettered = sum(1 for t in text_boxes if t and t[0] in 'ABCDE') >= 3
            has_lines = len(line_shapes) >= 3

            if has_numbered and has_lettered and has_lines:
                print(f"PASS: Component 7 — Slide 9 quiz: numbered items, lettered items, {len(line_shapes)} connecting lines (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — Slide 9 quiz: numbered={has_numbered}, lettered={has_lettered}, lines={len(line_shapes)} (need >=3)")
        else:
            print("FAIL: Component 7 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # ---- Component 8: Slide 10 gift shop items + museum info (0.05 pts) ----
    try:
        if num_slides >= 10:
            slide10 = prs.slides[9]
            all_text = " ".join(get_all_text(s) for s in slide10.shapes).lower()
            has_prices = '$' in all_text
            has_museum_info = 'museum' in all_text or 'hours' in all_text or 'am' in all_text
            has_gift_items = 'gift' in all_text or 'shop' in all_text or 'figurine' in all_text or 'jewelry' in all_text

            if has_prices and has_museum_info and has_gift_items:
                print(f"PASS: Component 8 — Slide 10 gift shop + museum info (0.05 pts)")
                total_score += 0.05
            else:
                print(f"FAIL: Component 8 — prices={has_prices}, museum_info={has_museum_info}, gift_items={has_gift_items}")
        else:
            print("FAIL: Component 8 — not enough slides")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # ---- Component 9: Colors #C9A959 and #263238 used (0.10 pts) ----
    try:
        colors = get_shape_colors(prs)
        has_gold = 'C9A959' in colors
        has_charcoal = '263238' in colors

        if has_gold and has_charcoal:
            print(f"PASS: Component 9 — Both colors found: C9A959, 263238 (0.10 pts)")
            total_score += 0.10
        elif has_gold or has_charcoal:
            partial = 0.05
            found = []
            if has_gold:
                found.append('C9A959')
            if has_charcoal:
                found.append('263238')
            print(f"PARTIAL: Component 9 — Only found: {found} ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 9 — Neither C9A959 nor 263238 found. Colors present: {colors}")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
