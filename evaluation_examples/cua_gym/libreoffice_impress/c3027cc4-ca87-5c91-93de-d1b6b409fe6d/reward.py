"""
FINAL REWARD SCRIPT - SUCCESS
Task: Remove outer borders but keep inner grid lines for Table 1.
Generated: 2025-10-17 09:12:14
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# -----------------------------------------------------------------------------
# Reward Script: Verify that outer borders are removed but inner grid lines are
#                retained for Table 1 in the presentation file.
# -----------------------------------------------------------------------------
# Scoring (progressive):
#   • 0.5 points – All OUTER borders of Table-1 are removed (no line definitions)
#   • 0.5 points – At least one INNER grid line is present (table still has
#                   internal borders)
#   => 1.0 only when both conditions hold
# -----------------------------------------------------------------------------
# NOTE:  Absolutely NO points are awarded for natural conditions such as the
#        mere existence of a presentation or a table. Points are awarded solely
#        for meeting the specific task requirements.
# -----------------------------------------------------------------------------

FILE_PATH = "/home/user/remove_outer_borders_but_keep_inner_grid_lines_for_table_1.pptx"

# Helper ----------------------------------------------------------------------

def _border_width(cell, side_tag):
    """Return border width (EMU) for a table cell side (L, R, T, B)."""
    tcPr = cell._tc.tcPr  # low-level XML for the table cell
    if tcPr is None:
        return 0
    ln = tcPr.find(
        f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}ln{side_tag}"
    )
    if ln is None:
        return 0
    w_attr = ln.get("w")
    try:
        return int(w_attr) if w_attr is not None else 0
    except ValueError:
        return 0

# Verification ----------------------------------------------------------------

def verify_remove_outer_borders_keep_inner(file_path: str) -> float:
    print(f"Verifying presentation file: {file_path}")

    if not os.path.exists(file_path):
        print("✗ File does not exist – cannot verify task")
        return 0.0

    # Load presentation -------------------------------------------------------
    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load PPTX: {e}")
        return 0.0

    # Locate the FIRST table (Table-1) ---------------------------------------
    table = None
    for s_idx, slide in enumerate(prs.slides, start=1):
        for sh_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_table:
                table = shape.table
                print(f"✓ Found Table 1 on slide {s_idx}, shape {sh_idx}")
                break
        if table is not None:
            break

    if table is None:
        print("✗ No table found – task not satisfied")
        return 0.0

    rows, cols = len(table.rows), len(table.columns)
    print(f"Table dimensions: {rows} rows × {cols} cols")

    # Flags ------------------------------------------------------------------
    outer_present = False   # any outer border exists?
    inner_present = False   # any inner grid line exists?

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)

            # --- OUTER borders ------------------------------------------------
            if r == 0        and _border_width(cell, "T") > 0:
                outer_present = True
            if r == rows-1   and _border_width(cell, "B") > 0:
                outer_present = True
            if c == 0        and _border_width(cell, "L") > 0:
                outer_present = True
            if c == cols-1   and _border_width(cell, "R") > 0:
                outer_present = True

            # --- INNER grid lines -------------------------------------------
            if r > 0 and _border_width(cell, "T") > 0:
                inner_present = True
            if c > 0 and _border_width(cell, "L") > 0:
                inner_present = True

    print(f"Outer borders present: {outer_present}")
    print(f"Inner grid lines present: {inner_present}")

    # Scoring ----------------------------------------------------------------
    score = 0.0

    if not outer_present:
        print("✓ Outer borders removed (0.5 points)")
        score += 0.5
    else:
        print("✗ Outer borders still present (0 points)")

    if inner_present:
        print("✓ Inner grid lines retained (0.5 points)")
        score += 0.5
    else:
        print("✗ Inner grid lines missing (0 points)")

    final_score = min(score, 1.0)
    print(f"Total score: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Execute verification and print REWARD ---------------------------------------
# -----------------------------------------------------------------------------

if __name__ == "__main__":
    reward = verify_remove_outer_borders_keep_inner(FILE_PATH)
    print(f"REWARD: {reward}")
