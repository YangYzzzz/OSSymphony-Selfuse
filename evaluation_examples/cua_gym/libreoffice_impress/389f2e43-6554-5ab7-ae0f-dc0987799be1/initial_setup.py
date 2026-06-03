"""
Initial Setup: Unhide all slides in presentation
Task ID: impress_fix_046
Domain: libreoffice_impress

Creates a 25-slide business presentation where slides 4, 8, 13, 17, and 22
(1-indexed) are hidden. The agent must unhide all slides.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'

# 1-indexed slide numbers that should be hidden
HIDDEN_SLIDES = {4, 8, 13, 17, 22}


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


# Slide content data — a realistic quarterly business review deck
SLIDE_DATA = [
    # (layout_idx, title, body_text)
    (0, "Q1 2025 Business Review", "Meridian Consulting Group\nQuarterly Performance Summary"),
    (1, "Agenda", "- Financial Overview\n- Client Engagement Metrics\n- Team Performance\n- Strategic Initiatives\n- Risk Assessment\n- Next Quarter Outlook"),
    (1, "Revenue Summary", "Total Revenue: $4.82M\nGross Margin: 42.3%\nOperating Income: $1.14M\nYear-over-Year Growth: 18.7%"),
    (1, "Regional Breakdown", "North America: $2.41M (+22%)\nEurope: $1.28M (+15%)\nAsia-Pacific: $0.74M (+11%)\nLatin America: $0.39M (+8%)"),
    (1, "Client Acquisition", "New Clients: 14\nTotal Active Clients: 87\nClient Retention Rate: 94.2%\nAverage Contract Value: $156K"),
    (1, "Top Client Accounts", "1. Axiom Technologies — $420K\n2. BluePeak Financial — $385K\n3. Crestline Manufacturing — $310K\n4. DataForge Analytics — $275K\n5. Evergreen Health Systems — $248K"),
    (1, "Service Line Performance", "Management Consulting: $1.92M\nIT Advisory: $1.45M\nHuman Capital: $0.88M\nFinancial Advisory: $0.57M"),
    (1, "Utilization Rates", "Overall Utilization: 78.4%\nSenior Consultants: 84.2%\nManagers: 81.7%\nAssociates: 72.1%\nTarget: 80%"),
    (1, "Employee Headcount", "Total Employees: 142\nNew Hires Q1: 18\nVoluntary Attrition: 3\nOpen Positions: 12"),
    (1, "Training & Development", "Training Hours Per Employee: 24\nCertifications Completed: 31\nLeadership Program Enrollees: 8\nMentor Pairs Active: 22"),
    (1, "Project Delivery Metrics", "Projects Completed: 28\nOn-Time Delivery: 89%\nClient Satisfaction Score: 4.6/5.0\nAverage Project Duration: 4.2 months"),
    (1, "Technology Investments", "Cloud Migration Budget: $320K\nCRM Platform Upgrade: $85K\nData Analytics Tools: $120K\nCybersecurity Enhancement: $95K"),
    (1, "Risk Register", "1. Talent shortage in AI/ML specialists — High\n2. Currency fluctuation impact on EU revenue — Medium\n3. Key client contract renewals pending — Medium\n4. Regulatory changes in data privacy — Low"),
    (1, "Competitive Landscape", "Market Position: Top 15 in mid-market consulting\nWin Rate on Proposals: 34%\nKey Differentiator: Industry-specific digital transformation\nCompetitor Activity: Increased M&A in sector"),
    (1, "Marketing & Brand", "Website Traffic: +28% QoQ\nLinkedIn Followers: 12,400 (+1,200)\nThought Leadership Articles: 8\nConference Presentations: 5\nMedia Mentions: 14"),
    (1, "Partnership Updates", "Microsoft Gold Partner — Renewed\nSalesforce Consulting Partner — New\nAWS Advanced Partner — In Progress\nSAP Silver Partner — Active"),
    (1, "Innovation Lab Initiatives", "Generative AI Pilot: 3 client engagements\nProcess Mining Tool: Beta testing with 2 clients\nESG Analytics Dashboard: Development phase\nPredictive Staffing Model: Internal rollout"),
    (1, "Financial Forecast Q2", "Projected Revenue: $5.15M\nTarget Margin: 44%\nPlanned Investments: $180K\nExpected Headcount: 148"),
    (1, "Strategic Priorities Q2", "1. Launch AI advisory practice\n2. Expand European presence (Berlin office)\n3. Achieve 82% utilization target\n4. Complete CRM migration\n5. Hire 3 senior partners"),
    (1, "Client Success Story", "Axiom Technologies Digital Transformation\n- Duration: 8 months\n- Team: 12 consultants\n- Outcome: 35% operational efficiency gain\n- Contract extension: 2 additional phases"),
    (1, "Sustainability Report", "Carbon Offset: 120 tons\nRemote Work Ratio: 62%\nPaper Reduction: 40% vs Q1 2024\nGreen Office Certifications: 3 of 5 locations"),
    (1, "Board Recommendations", "1. Approve Berlin office expansion ($450K)\n2. Authorize AI practice hiring (5 FTEs)\n3. Increase training budget by 15%\n4. Renew cybersecurity insurance"),
    (1, "Key Performance Indicators", "Revenue Growth: 18.7% (Target: 15%) ✓\nClient Retention: 94.2% (Target: 90%) ✓\nUtilization: 78.4% (Target: 80%) ✗\nEmployee Satisfaction: 4.1/5 (Target: 4.0) ✓"),
    (1, "Appendix: Financial Details", "Detailed P&L available in supplementary document\nBalance Sheet: Net assets $8.2M\nCash Position: $2.1M\nAccounts Receivable: $1.8M (DSO: 42 days)"),
    (0, "Thank You", "Questions & Discussion\nContact: strategy@meridiancg.com"),
]


def create_initial():
    prs = Presentation()

    for idx, (layout_idx, title, body) in enumerate(SLIDE_DATA):
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(28) if layout_idx == 0 else Pt(24)
                run.font.bold = True

        # Set body content
        if layout_idx == 0:
            # Title slide — use subtitle placeholder
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = body
                for run in slide.placeholders[1].text_frame.paragraphs[0].runs:
                    run.font.size = Pt(18)
        elif layout_idx == 1:
            # Content slide — use body placeholder
            if len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                lines = body.split('\n')
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.space_after = Pt(6)
                    for run in p.runs:
                        run.font.size = Pt(16)

        # Hide designated slides (1-indexed: slide_num = idx + 1)
        slide_num = idx + 1
        if slide_num in HIDDEN_SLIDES:
            slide._element.set('show', '0')

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')
    print(f'Hidden slides (1-indexed): {sorted(HIDDEN_SLIDES)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
