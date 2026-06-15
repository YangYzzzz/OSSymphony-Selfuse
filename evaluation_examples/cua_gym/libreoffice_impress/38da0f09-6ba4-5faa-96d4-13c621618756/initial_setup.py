"""
Initial Setup: Create thesis presentation with 15 slides and university logo file.
Task ID: impress_stu_038
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
LOGO_DIR = f'{WORKDIR}/Downloads'
LOGO_PATH = f'{LOGO_DIR}/university_logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a simple university logo PNG."""
    os.makedirs(LOGO_DIR, exist_ok=True)
    img = Image.new('RGBA', (200, 200), (0, 0, 0, 0))
    from PIL import ImageDraw, ImageFont
    draw = ImageDraw.Draw(img)
    # Shield shape background
    draw.ellipse([10, 10, 190, 190], fill=(0, 51, 102, 255))
    # Inner circle
    draw.ellipse([30, 30, 170, 170], fill=(255, 255, 255, 255))
    # Inner blue circle
    draw.ellipse([50, 50, 150, 150], fill=(0, 51, 102, 255))
    # Letter U
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 60)
    except:
        font = ImageFont.load_default()
    draw.text((75, 65), "U", fill=(255, 255, 255, 255), font=font)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def add_text_to_placeholder(shape, text, font_size=Pt(18), bold=False, color=None):
    """Helper to set text on a placeholder shape."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=Pt(14),
                bold=False, color=None, alignment=None):
    """Add a textbox with text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    for run in p.runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=Pt(14)):
    """Add a textbox with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = font_size
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Machine Learning Approaches for\nUrban Traffic Flow Prediction"
    slide.placeholders[1].text = "Elena Vasquez\nDepartment of Computer Science\nStanford University\nThesis Defense - March 2026"

    # --- Slide 2: Outline ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Presentation Outline", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "1. Introduction and Motivation",
        "2. Literature Review",
        "3. Research Questions and Hypotheses",
        "4. Methodology",
        "5. Data Collection and Preprocessing",
        "6. Model Architecture",
        "7. Experimental Results",
        "8. Discussion and Analysis",
        "9. Conclusions and Future Work",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(5), items, Pt(16))

    # --- Slide 3: Introduction ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Introduction", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
                "Urban traffic congestion costs the U.S. economy an estimated $87 billion "
                "annually in lost productivity and wasted fuel. Accurate short-term traffic flow "
                "prediction enables dynamic signal control, route guidance, and infrastructure "
                "planning. Traditional statistical methods (ARIMA, VAR) struggle to capture the "
                "complex spatiotemporal dependencies inherent in urban traffic networks.",
                Pt(14))

    # --- Slide 4: Motivation ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Motivation", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "Growing urban populations increase traffic complexity",
        "Real-time prediction needed for autonomous vehicle integration",
        "Existing models fail during non-recurring congestion events",
        "Limited research on multi-modal transportation networks",
        "Need for interpretable predictions for urban planners",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(5), items, Pt(15))

    # --- Slide 5: Literature Review ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Literature Review", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
                "Key prior work:\n\n"
                "- Li et al. (2018): Diffusion Convolutional RNN for traffic forecasting\n"
                "- Yu et al. (2019): Spatio-Temporal Graph Convolutional Networks (STGCN)\n"
                "- Zheng et al. (2020): GMAN attention-based model, METR-LA dataset\n"
                "- Wu et al. (2021): Graph WaveNet with adaptive adjacency matrices\n"
                "- Song et al. (2022): Spatial-Temporal Synchronous Graph Conv Networks\n\n"
                "Gap: None address multi-modal transit (bus, subway, bike-share) jointly.",
                Pt(13))

    # --- Slide 6: Research Questions ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Research Questions", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "RQ1: Can graph neural networks effectively model multi-modal urban transit?",
        "RQ2: Does incorporating weather and event data improve prediction accuracy?",
        "RQ3: How does prediction horizon affect model performance (5 vs 15 vs 30 min)?",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(4.5), items, Pt(15))

    # --- Slide 7: Methodology Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Methodology Overview", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
                "Our approach, Multi-Modal Spatio-Temporal Graph Network (MMSTGN), consists of:\n\n"
                "1. Dynamic graph construction from GPS trajectory data\n"
                "2. Multi-modal feature fusion layer combining traffic, transit, and bike-share\n"
                "3. Temporal attention mechanism for capturing periodic patterns\n"
                "4. Graph convolutional layers with adaptive adjacency learning\n"
                "5. Multi-horizon prediction head with uncertainty estimation",
                Pt(14))

    # --- Slide 8: Data Collection ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Data Collection", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    # Add a table
    rows, cols = 5, 4
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(8.2), Inches(3))
    table = tbl_shape.table
    headers = ["Data Source", "Records", "Time Period", "Granularity"]
    data = [
        ["Loop Detectors (PeMS)", "2.4M", "Jan 2023 - Dec 2024", "5 min"],
        ["Bus GPS Traces", "856K", "Jan 2023 - Dec 2024", "30 sec"],
        ["Subway Turnstile", "12.1M", "Jan 2023 - Dec 2024", "15 min"],
        ["Bike-Share Trips", "3.7M", "Jan 2023 - Dec 2024", "Trip-level"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 9: Data Preprocessing ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Data Preprocessing", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "Missing value imputation using spatiotemporal kriging",
        "Outlier detection via Isolation Forest (removed 2.3% of records)",
        "Temporal alignment across data sources to 5-minute intervals",
        "Z-score normalization per sensor with rolling 7-day statistics",
        "Graph construction: 325 nodes, 1,847 edges (distance threshold 5km)",
        "70/15/15 train-validation-test split (chronological)",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(5), items, Pt(14))

    # --- Slide 10: Model Architecture ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Model Architecture", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
                "MMSTGN Architecture Details:\n\n"
                "- Input: 12 historical time steps (1 hour) per node\n"
                "- Feature dimension: 64 per modality, 192 fused\n"
                "- Graph convolution: 3 layers, Chebyshev order K=3\n"
                "- Temporal attention: 8 heads, 256 hidden units\n"
                "- Output: 1/3/6 steps ahead (5/15/30 minutes)\n"
                "- Parameters: 1.2M trainable, 45 epochs, Adam optimizer\n"
                "- Learning rate: 1e-3 with cosine annealing schedule",
                Pt(14))

    # --- Slide 11: Results - Quantitative ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Experimental Results", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    rows, cols = 5, 4
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(8.2), Inches(3.5))
    table = tbl_shape.table
    headers = ["Model", "MAE", "RMSE", "MAPE (%)"]
    data = [
        ["ARIMA", "7.42", "12.83", "18.6"],
        ["STGCN", "4.31", "8.17", "10.2"],
        ["Graph WaveNet", "3.89", "7.52", "9.1"],
        ["MMSTGN (Ours)", "3.24", "6.41", "7.8"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # --- Slide 12: Results - Analysis ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Results Analysis", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "MMSTGN achieves 16.7% MAE reduction over Graph WaveNet",
        "Multi-modal fusion contributes 8.3% of the improvement",
        "Performance gap widens at longer prediction horizons",
        "Largest gains during peak hours and incident scenarios",
        "Attention weights reveal interpretable spatial patterns",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(5), items, Pt(15))

    # --- Slide 13: Discussion ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Discussion", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.5), Inches(5),
                "Key findings support all three research hypotheses. The multi-modal "
                "approach significantly outperforms single-source baselines, particularly "
                "during non-recurring congestion events where subway and bus data provide "
                "complementary signals.\n\n"
                "Limitations:\n"
                "- Computational cost 3x higher than Graph WaveNet\n"
                "- Requires all four data modalities for best performance\n"
                "- Tested only on San Francisco Bay Area network",
                Pt(14))

    # --- Slide 14: Conclusions ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Conclusions", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "Novel MMSTGN architecture for multi-modal traffic prediction",
        "State-of-the-art results on PeMS-BAY benchmark (MAE: 3.24)",
        "First model to jointly incorporate bus, subway, and bike-share data",
        "Interpretable attention mechanism for urban planning applications",
        "Open-sourced code and pre-trained models for reproducibility",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(5), items, Pt(15))

    # --- Slide 15: Future Work & Questions ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(1),
                "Future Work", Pt(28), bold=True, color=RGBColor(0x00, 0x33, 0x66))
    items = [
        "Extend to additional cities (NYC, Chicago, London)",
        "Incorporate ride-sharing and e-scooter data modalities",
        "Real-time deployment with streaming graph updates",
        "Integration with traffic signal control optimization",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(3.5), items, Pt(15))
    add_textbox(slide, Inches(2), Inches(5.5), Inches(6), Inches(1.5),
                "Thank You\nQuestions?", Pt(24), bold=True,
                color=RGBColor(0x00, 0x33, 0x66), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch GUI
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create the logo first, then the presentation
create_logo()
create_initial()
