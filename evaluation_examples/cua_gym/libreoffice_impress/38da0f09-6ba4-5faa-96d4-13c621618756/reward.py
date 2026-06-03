"""
Reward Script: Set up master slide formatting with university logo and bottom line
Task ID: impress_stu_038
Domain: libreoffice_impress
Scoring:
  Component 1: Logo image on slide master (0.3)
  Component 2: Logo position top-right, ~1 inch square (0.3)
  Component 3: Horizontal line/rectangle near bottom of slide master (0.2)
  Component 4: Line is dark gray #333333 and thin (0.2)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_038'

# Tolerances
INCH = 914400  # EMU per inch
POS_TOLERANCE = 0.15  # 15% relative tolerance for position checks


def is_approx(val, expected, tolerance=POS_TOLERANCE):
    """Check if val is approximately equal to expected within relative tolerance."""
    if expected == 0:
        return abs(val) < INCH * 0.2  # within 0.2 inch of zero
    return abs(val - expected) / abs(expected) <= tolerance


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_width = prs.slide_width   # 9144000 EMU (10 inches)
    slide_height = prs.slide_height  # 6858000 EMU (7.5 inches)

    # Access the first slide master
    try:
        master = prs.slide_masters[0]
    except Exception as e:
        print(f"CRITICAL: Cannot access slide master: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find picture and line shapes on the master
    logo_shape = None
    line_shape = None

    for shape in master.shapes:
        # Find picture (logo)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            logo_shape = shape
        # Find thin rectangle (line) - height should be small (thin)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # A "line" implemented as a thin rectangle: height < 0.25 inch
            if shape.height < INCH * 0.25:
                line_shape = shape

    # Component 1: Logo image exists on slide master (0.3 points)
    try:
        if logo_shape is not None:
            # Verify it's actually an image with content
            blob_size = len(logo_shape.image.blob)
            if blob_size > 0:
                print(f"PASS: Component 1 -- Logo image found on slide master, blob={blob_size} bytes (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 1 -- Logo image has empty blob")
        else:
            print(f"FAIL: Component 1 -- No picture shape found on slide master")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Logo positioned in top-right corner, approximately 1 inch square (0.3 points)
    try:
        if logo_shape is not None:
            left = logo_shape.left
            top = logo_shape.top
            width = logo_shape.width
            height = logo_shape.height
            right_edge = left + width

            print(f"  Logo position: left={left}, top={top}, width={width}, height={height}")
            print(f"  Slide dimensions: width={slide_width}, height={slide_height}")
            print(f"  Logo right edge: {right_edge}, slide right: {slide_width}")

            # Top-right means: right edge near slide right edge, top near 0
            # ~1 inch square means width and height each ~ 914400 EMU
            points_earned = 0.0

            # Check top-right positioning (logo right edge within ~1.5 inches of slide right)
            distance_from_right = slide_width - right_edge
            is_right = distance_from_right < INCH * 1.5
            is_top = top < INCH * 1.0  # top within 1 inch from top edge

            if is_right and is_top:
                points_earned += 0.15
                print(f"  PASS: Logo is in top-right area (dist_from_right={distance_from_right/INCH:.2f}in, top={top/INCH:.2f}in)")
            else:
                print(f"  FAIL: Logo not in top-right (dist_from_right={distance_from_right/INCH:.2f}in, top={top/INCH:.2f}in)")

            # Check ~1 inch square (within 50% tolerance since task says "approximately")
            is_1inch_w = abs(width - INCH) / INCH <= 0.5
            is_1inch_h = abs(height - INCH) / INCH <= 0.5

            if is_1inch_w and is_1inch_h:
                points_earned += 0.15
                print(f"  PASS: Logo is approximately 1 inch square (w={width/INCH:.2f}in, h={height/INCH:.2f}in)")
            else:
                print(f"  FAIL: Logo not ~1 inch square (w={width/INCH:.2f}in, h={height/INCH:.2f}in)")

            if points_earned > 0:
                total_score += points_earned
                print(f"PASS: Component 2 -- Logo positioning ({points_earned} pts)")
            else:
                print(f"FAIL: Component 2 -- Logo positioning incorrect")
        else:
            print(f"FAIL: Component 2 -- No logo found, cannot check position")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Horizontal line/thin rectangle near bottom of master (0.2 points)
    try:
        if line_shape is not None:
            line_top = line_shape.top
            line_width = line_shape.width
            line_height = line_shape.height

            print(f"  Line shape: top={line_top}, width={line_width}, height={line_height}")

            # Must be near the bottom of the slide (within bottom 15% of slide)
            bottom_threshold = slide_height * 0.85
            is_near_bottom = line_top >= bottom_threshold

            # Must span a reasonable portion of the slide width (at least 50%)
            is_wide = line_width >= slide_width * 0.5

            if is_near_bottom and is_wide:
                print(f"PASS: Component 3 -- Horizontal line near bottom (top={line_top/INCH:.2f}in, width={line_width/INCH:.2f}in) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 -- Line not at bottom or not wide enough (near_bottom={is_near_bottom}, wide={is_wide})")
        else:
            print(f"FAIL: Component 3 -- No thin rectangular shape found on slide master")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Line is dark gray #333333 and thin (0.2 points)
    try:
        if line_shape is not None:
            line_height = line_shape.height
            is_thin = line_height <= INCH * 0.1  # thinner than 0.1 inch

            # Check fill color - derive color_ok from actual check
            dark_grays = {'333333', '3C3C3C', '2D2D2D', '404040', '383838'}
            actual_color = None

            try:
                fill = line_shape.fill
                if fill.type is not None:
                    actual_color = str(fill.fore_color.rgb)
            except Exception:
                pass

            # Fallback: check line outline color
            if actual_color is None:
                try:
                    ln = line_shape.line
                    if ln.fill.type is not None:
                        actual_color = str(ln.color.rgb)
                except Exception:
                    pass

            color_ok = actual_color in dark_grays if actual_color else False

            points_earned = 0.0
            if is_thin:
                points_earned += 0.1
                print(f"  PASS: Line is thin (height={line_height/INCH:.3f}in)")
            else:
                print(f"  FAIL: Line not thin (height={line_height/INCH:.3f}in)")

            if color_ok:
                points_earned += 0.1
                print(f"  PASS: Line color is dark gray ({actual_color})")
            else:
                print(f"  FAIL: Line color not dark gray #333333 (actual={actual_color})")

            if points_earned > 0:
                total_score += points_earned
                print(f"PASS: Component 4 -- Line style ({points_earned} pts)")
            else:
                print(f"FAIL: Component 4 -- Line style incorrect")
        else:
            print(f"FAIL: Component 4 -- No line shape found, cannot check style")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state(domain):
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state("libreoffice_impress")
    verify_task(file_path)
