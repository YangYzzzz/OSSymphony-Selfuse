"""
Reward Script: Add logo to Master Slide in bottom-left corner at 1.5cm x 1.5cm
Task ID: impress_gf5_003
Domain: libreoffice_impress
Scoring:
  Component 1 (0.3): Master slide contains a PICTURE shape (image on master)
  Component 2 (0.3): The image blob matches /home/user/assets/logo.png
  Component 3 (0.2): Image dimensions are 1.5cm x 1.5cm (540000 EMU each)
  Component 4 (0.2): Image is positioned in the bottom-left corner of the slide
"""

import os
import hashlib

from pptx import Presentation
from pptx.util import Cm, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_003'

# Tolerance for position/size checks (5% relative)
TOLERANCE = 0.05


def is_approx(val, expected, tol=TOLERANCE):
    """Check if val is approximately equal to expected within relative tolerance."""
    if expected == 0:
        return abs(val) < 10000  # ~0.03cm absolute tolerance for zero
    return abs(val - expected) / abs(expected) <= tol


def persist_app_state(domain):
    """Send Ctrl+S to save any unsaved LibreOffice edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_logo_md5():
    """Get MD5 hash of the source logo file."""
    logo_path = os.path.join(WORKDIR, 'assets', 'logo.png')
    if not os.path.exists(logo_path):
        return None
    with open(logo_path, 'rb') as f:
        return hashlib.md5(f.read()).hexdigest()


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Find all PICTURE shapes on the first slide master
    master = prs.slide_masters[0]
    picture_shapes = [s for s in master.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Component 1: Master slide contains a PICTURE shape (0.3 points)
    try:
        if len(picture_shapes) > 0:
            print(f"PASS: Component 1 - Master slide has {len(picture_shapes)} picture shape(s) (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 - No picture shapes found on master slide")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: The image blob matches the source logo.png (0.3 points)
    try:
        logo_md5 = get_logo_md5()
        if logo_md5 is None:
            print(f"FAIL: Component 2 - Logo file not found at /home/user/assets/logo.png")
        elif len(picture_shapes) == 0:
            print(f"FAIL: Component 2 - No picture on master to compare")
        else:
            # Check if any picture on master matches the logo
            matched = False
            for pic in picture_shapes:
                blob_md5 = hashlib.md5(pic.image.blob).hexdigest()
                if blob_md5 == logo_md5:
                    matched = True
                    break
            if matched:
                print(f"PASS: Component 2 - Master image matches logo.png (MD5: {logo_md5}) (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 - Master image does not match logo.png")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Image size is 1.5cm x 1.5cm (0.2 points)
    # 1.5cm = 540000 EMU (1cm = 360000 EMU)
    expected_size = 540000  # 1.5cm in EMU
    try:
        if len(picture_shapes) == 0:
            print(f"FAIL: Component 3 - No picture on master to check size")
        else:
            # Find the logo picture (matched by blob or first picture)
            logo_md5 = get_logo_md5()
            target_pic = None
            for pic in picture_shapes:
                blob_md5 = hashlib.md5(pic.image.blob).hexdigest()
                if blob_md5 == logo_md5:
                    target_pic = pic
                    break
            if target_pic is None:
                target_pic = picture_shapes[0]

            w = target_pic.width
            h = target_pic.height
            w_ok = is_approx(w, expected_size)
            h_ok = is_approx(h, expected_size)
            if w_ok and h_ok:
                print(f"PASS: Component 3 - Logo size {w/360000:.2f}cm x {h/360000:.2f}cm matches 1.5x1.5cm (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 3 - Logo size {w/360000:.2f}cm x {h/360000:.2f}cm, expected 1.5x1.5cm")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Image is in the bottom-left corner (0.2 points)
    # Bottom-left means: left edge near 0, bottom edge near slide_height
    # We check: left < 20% of slide width, and bottom edge > 80% of slide height
    try:
        if len(picture_shapes) == 0:
            print(f"FAIL: Component 4 - No picture on master to check position")
        else:
            logo_md5 = get_logo_md5()
            target_pic = None
            for pic in picture_shapes:
                blob_md5 = hashlib.md5(pic.image.blob).hexdigest()
                if blob_md5 == logo_md5:
                    target_pic = pic
                    break
            if target_pic is None:
                target_pic = picture_shapes[0]

            left = target_pic.left
            top = target_pic.top
            bottom_edge = top + target_pic.height
            right_edge = left + target_pic.width

            # Bottom-left: left edge in left 20% of slide, bottom edge in bottom 20% of slide
            in_left = left < slide_width * 0.20
            in_bottom = bottom_edge > slide_height * 0.80

            if in_left and in_bottom:
                print(f"PASS: Component 4 - Logo at bottom-left (left={left/360000:.2f}cm, bottom={bottom_edge/360000:.2f}cm) (0.2 pts)")
                total_score += 0.2
            else:
                print(f"FAIL: Component 4 - Logo not in bottom-left. left={left/360000:.2f}cm (need <{slide_width*0.20/360000:.2f}cm), bottom={bottom_edge/360000:.2f}cm (need >{slide_height*0.80/360000:.2f}cm)")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
