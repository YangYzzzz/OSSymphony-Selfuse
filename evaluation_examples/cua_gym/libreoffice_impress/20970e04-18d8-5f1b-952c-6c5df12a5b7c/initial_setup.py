"""
Initial Setup: Add logo to Master Slide in company_profile.pptx
Task ID: impress_gf5_003
Domain: libreoffice_impress

Creates a 12-slide company profile presentation and a 200x200 logo image.
No logo on any slide or master. Opens in LibreOffice Impress.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_003'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
ASSETS_DIR = f'{WORKDIR}/assets'
LOGO_PATH = f'{ASSETS_DIR}/logo.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_logo():
    """Create a 200x200 square logo PNG."""
    os.makedirs(ASSETS_DIR, exist_ok=True)
    img = Image.new('RGBA', (200, 200), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    # Blue circle background
    draw.ellipse([10, 10, 190, 190], fill=(0, 102, 204, 255))
    # White "A" letter in center (for "Apex Corp")
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 100)
    except (IOError, OSError):
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), "A", font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    x = (200 - tw) // 2
    y = (200 - th) // 2 - 10
    draw.text((x, y), "A", fill=(255, 255, 255, 255), font=font)
    img.save(LOGO_PATH)
    print(f'Logo created: {LOGO_PATH}')


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)
    return slide


def add_blank_slide_with_text(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title textbox
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    # Body textbox
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body_text
    p2.font.size = Pt(16)
    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Apex Corporation", "Company Profile 2025")

    # Slide 2: About Us
    add_content_slide(prs, "About Apex Corporation", [
        "Founded in 2008 in San Francisco, California",
        "Leading provider of enterprise cloud infrastructure",
        "Over 3,200 employees across 14 global offices",
        "Annual revenue exceeding $890 million (FY2024)",
        "Recognized as a Gartner Magic Quadrant Leader",
    ])

    # Slide 3: Mission & Vision
    add_content_slide(prs, "Mission & Vision", [
        "Mission: Empowering businesses with scalable, secure cloud solutions",
        "Vision: A world where every organization can innovate without infrastructure limits",
        "Core Values: Innovation, Integrity, Customer Success, Collaboration",
    ])

    # Slide 4: Leadership Team
    add_content_slide(prs, "Leadership Team", [
        "CEO: Dr. Priya Sharma - Former VP Engineering at CloudScale",
        "CTO: James Whitfield - 20 years in distributed systems",
        "CFO: Maria Gonzalez - Led IPO of TechVenture in 2019",
        "COO: David Park - Operations expertise from Fortune 500",
        "CMO: Rachel Thompson - Brand strategist, ex-Salesforce",
    ])

    # Slide 5: Products & Services
    add_content_slide(prs, "Products & Services", [
        "ApexCloud Platform - Multi-cloud management and orchestration",
        "ApexSecure - Zero-trust security framework for hybrid environments",
        "ApexAnalytics - Real-time data pipeline and visualization suite",
        "ApexEdge - Edge computing solutions for IoT deployments",
        "Professional Services - Migration, training, and 24/7 support",
    ])

    # Slide 6: Market Presence
    add_blank_slide_with_text(prs, "Global Market Presence",
        "Apex Corporation operates in 28 countries with primary hubs in "
        "San Francisco, London, Singapore, and Tokyo. Our customer base includes "
        "over 4,500 enterprise clients spanning financial services, healthcare, "
        "retail, and government sectors. In 2024, we expanded into Latin America "
        "with new offices in Sao Paulo and Mexico City.")

    # Slide 7: Financial Highlights
    add_content_slide(prs, "Financial Highlights (FY2024)", [
        "Revenue: $890M (+22% YoY)",
        "Gross Margin: 72.4%",
        "Net Income: $134M",
        "Free Cash Flow: $198M",
        "R&D Investment: $210M (24% of revenue)",
        "Customer Retention Rate: 96.8%",
    ])

    # Slide 8: Technology Stack
    add_content_slide(prs, "Technology & Innovation", [
        "Kubernetes-native architecture across all products",
        "Proprietary AI/ML engine for workload optimization",
        "200+ patents in cloud computing and security",
        "Open-source contributor: 45 active projects on GitHub",
        "SOC 2 Type II, ISO 27001, HIPAA compliant",
    ])

    # Slide 9: Customer Success Stories
    add_content_slide(prs, "Customer Success Stories", [
        "GlobalBank: Reduced infrastructure costs by 42% in 18 months",
        "MedTech Solutions: Achieved HIPAA compliance in 6 weeks",
        "RetailMax: Scaled to handle 10x traffic during Black Friday",
        "City of Portland: Modernized citizen services platform",
        "Quantum Aerospace: Real-time telemetry processing at the edge",
    ])

    # Slide 10: Partnerships
    add_content_slide(prs, "Strategic Partnerships", [
        "AWS Advanced Technology Partner",
        "Microsoft Azure Gold Partner",
        "Google Cloud Premier Partner",
        "VMware Principal Partner",
        "Certified integrations with ServiceNow, Splunk, and Datadog",
    ])

    # Slide 11: Sustainability
    add_blank_slide_with_text(prs, "Sustainability Commitment",
        "Apex Corporation is committed to achieving carbon neutrality by 2027. "
        "Our data centers run on 85% renewable energy, and we have reduced our "
        "carbon footprint by 38% since 2020. We participate in the Science Based "
        "Targets initiative and publish an annual ESG report. In 2024, we planted "
        "over 50,000 trees through our partnership with One Tree Planted.")

    # Slide 12: Contact
    add_content_slide(prs, "Contact Us", [
        "Headquarters: 450 Mission Street, San Francisco, CA 94105",
        "Phone: +1 (415) 555-0200",
        "Email: info@apexcorp.com",
        "Website: www.apexcorporation.com",
        "LinkedIn: linkedin.com/company/apex-corporation",
    ])

    prs.save(OUTPUT)
    print(f'Initial presentation created: {OUTPUT}')


def main():
    create_logo()
    create_presentation()
    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


main()
