"""
Reward Script: Annotated bibliography slide with APA formatting, sidebar, grouped references
Task ID: impress_stu_095
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Colored sidebar rectangle (#2C3E50, 0.5" wide, full height)
  Component 2 (0.25): Three bold section headers (Journal Articles, Books, Online Sources)
  Component 3 (0.25): Six references in Calibri ~11pt
  Component 4 (0.25): Hanging indent formatting on references
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_095'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 9 slides
    if len(prs.slides) < 9:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 9")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[8]  # slide 9, 0-indexed

    # =========================================================================
    # Component 1: Colored sidebar rectangle (0.25 points)
    # Must be a rectangle shape on the left, ~0.5 inches wide, full slide height,
    # filled with #2C3E50.
    # This shape does NOT exist in initial_env (initial has only 2 shapes).
    # =========================================================================
    try:
        sidebar_score = 0.0
        slide_height = prs.slide_height
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check position: should be at or near left=0, top=0
                # Width should be ~0.5 inches (457200 EMU), height ~full slide height
                width_ok = abs(shape.width - Inches(0.5)) < Inches(0.15)
                height_ok = abs(shape.height - slide_height) < Inches(0.5)
                left_ok = shape.left < Inches(0.25)

                if width_ok and height_ok and left_ok:
                    # Check fill color
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID
                        color_rgb = str(fill.fore_color.rgb).upper()
                        if color_rgb == '2C3E50':
                            sidebar_score = 0.25
                            print(f"PASS: Component 1 -- Sidebar found: {shape.width}x{shape.height} EMU, color #{color_rgb} (0.25 pts)")
                            total_score += 0.25
                            break
                        else:
                            print(f"FAIL: Component 1 -- Sidebar shape found but color is #{color_rgb}, expected #2C3E50")
                    else:
                        print(f"FAIL: Component 1 -- Sidebar shape found but fill type is {fill.type}, expected SOLID (1)")

        if sidebar_score < 0.01:
            print("FAIL: Component 1 -- No sidebar rectangle found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================================
    # Component 2: Three bold section headers (0.25 points)
    # Headers: "Journal Articles", "Books", "Online Sources"
    # These do NOT exist in initial_env.
    # =========================================================================
    try:
        expected_headers = {"journal articles", "books", "online sources"}
        found_headers = set()

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip().lower()
                    if text in expected_headers:
                        # Check if bold
                        runs = [r for r in para.runs if (r.text or "").strip()]
                        if runs and runs[0].font.bold:
                            found_headers.add(text)

        if len(found_headers) == 3:
            print(f"PASS: Component 2 -- All 3 bold headers found: {found_headers} (0.25 pts)")
            total_score += 0.25
        else:
            missing = expected_headers - found_headers
            print(f"FAIL: Component 2 -- Found {len(found_headers)}/3 headers. Missing: {missing}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================================
    # Component 3: Six references in Calibri ~11pt (0.25 points)
    # References are non-header, non-empty paragraphs with author-style text.
    # initial_env has NO references at all. Score only if >= 6 refs found.
    # 11pt = 139700 EMU (size). Calibri font.
    # =========================================================================
    try:
        header_texts = {"journal articles", "books", "online sources", "annotated bibliography"}
        reference_count = 0
        calibri_11pt_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if text.lower() in header_texts:
                        continue
                    # This is a reference paragraph
                    reference_count += 1
                    runs = [r for r in para.runs if (r.text or "").strip()]
                    if runs:
                        r = runs[0]
                        font_name = r.font.name
                        font_size = r.font.size
                        is_calibri = font_name is not None and font_name.lower() == 'calibri'
                        # 11pt = 139700 EMU, allow small tolerance
                        is_11pt = font_size is not None and abs(font_size - 139700) < 15000
                        if is_calibri and is_11pt:
                            calibri_11pt_count += 1

        if reference_count >= 6 and calibri_11pt_count >= 6:
            print(f"PASS: Component 3 -- {reference_count} references found, {calibri_11pt_count} in Calibri ~11pt (0.25 pts)")
            total_score += 0.25
        elif reference_count >= 6:
            print(f"FAIL: Component 3 -- {reference_count} references found but only {calibri_11pt_count} in Calibri ~11pt")
        else:
            print(f"FAIL: Component 3 -- Only {reference_count} references found (need >= 6)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # =========================================================================
    # Component 4: Hanging indent formatting on references (0.25 points)
    # marL=457200 (0.5 inches), indent=-457200 (-0.5 inches) for APA hanging indent.
    # initial_env has NO such formatting.
    # =========================================================================
    try:
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        hanging_indent_count = 0
        total_refs = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if text.lower() in header_texts:
                        continue
                    # Reference paragraph
                    total_refs += 1
                    pPr = para._p.find(f'{ns_a}pPr')
                    if pPr is not None:
                        marL = pPr.get('marL')
                        indent = pPr.get('indent')
                        if marL is not None and indent is not None:
                            marL_val = int(marL)
                            indent_val = int(indent)
                            # Hanging indent: positive marL, negative indent
                            if marL_val > 0 and indent_val < 0:
                                hanging_indent_count += 1

        if total_refs >= 6 and hanging_indent_count >= 6:
            print(f"PASS: Component 4 -- {hanging_indent_count}/{total_refs} references have hanging indent (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 -- {hanging_indent_count}/{total_refs} references have hanging indent (need >= 6)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
