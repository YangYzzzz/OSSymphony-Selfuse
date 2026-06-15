"""
Initial Setup: Create a cognitive psychology presentation with 9 slides.
Slide 9 is titled 'Annotated Bibliography' but is otherwise empty.
Task ID: impress_stu_095
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_095'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Cognitive Psychology: A Comprehensive Review",
        "Department of Psychology | Spring 2026 Seminar"
    )

    # Slide 2: Introduction
    add_content_slide(prs, "Introduction to Cognitive Psychology", [
        "Study of mental processes including perception, memory, and reasoning",
        "Emerged as a distinct field in the late 1950s",
        "Combines experimental methods with computational modeling",
        "Key contributors: Ulric Neisser, George Miller, Noam Chomsky",
        "Interdisciplinary connections with neuroscience and AI"
    ])

    # Slide 3: Memory Systems
    add_content_slide(prs, "Memory Systems and Processes", [
        "Sensory memory: iconic (visual) and echoic (auditory) stores",
        "Short-term / working memory: capacity of 7 +/- 2 items (Miller, 1956)",
        "Long-term memory: declarative (explicit) vs. procedural (implicit)",
        "Encoding specificity principle (Tulving & Thomson, 1973)",
        "Consolidation and reconsolidation during sleep"
    ])

    # Slide 4: Attention
    add_content_slide(prs, "Attention and Cognitive Control", [
        "Selective attention: Broadbent's filter model vs. late selection",
        "Divided attention and the cocktail party effect",
        "Inattentional blindness (Simons & Chabris, 1999)",
        "Executive functions and the prefrontal cortex",
        "Attentional blink and temporal processing limits"
    ])

    # Slide 5: Language Processing
    add_content_slide(prs, "Language and Communication", [
        "Phonological processing and speech perception",
        "Syntactic parsing and garden-path sentences",
        "Semantic networks and spreading activation models",
        "Bilingual language processing and code-switching",
        "Language production: from conceptualization to articulation"
    ])

    # Slide 6: Decision Making
    add_content_slide(prs, "Decision Making and Reasoning", [
        "Dual-process theory: System 1 (fast/intuitive) vs. System 2 (slow/deliberate)",
        "Heuristics and biases framework (Tversky & Kahneman, 1974)",
        "Prospect theory and loss aversion",
        "Bounded rationality and satisficing (Simon, 1957)",
        "Ecological rationality and adaptive toolbox"
    ])

    # Slide 7: Cognitive Development
    add_content_slide(prs, "Cognitive Development Across the Lifespan", [
        "Piaget's stages: sensorimotor through formal operational",
        "Vygotsky's zone of proximal development and scaffolding",
        "Theory of mind development in early childhood",
        "Cognitive reserve and aging: protective factors",
        "Neuroplasticity and lifelong learning potential"
    ])

    # Slide 8: Research Methods
    add_content_slide(prs, "Research Methods in Cognitive Psychology", [
        "Behavioral paradigms: reaction time, accuracy, priming effects",
        "Neuroimaging: fMRI, EEG, MEG, and PET methodologies",
        "Computational modeling: connectionist and Bayesian approaches",
        "Eye-tracking and pupillometry for real-time processing",
        "Ecological validity and naturalistic study designs"
    ])

    # Slide 9: Annotated Bibliography - EMPTY (title only)
    slide9 = add_title_only_slide(prs, "Annotated Bibliography")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
