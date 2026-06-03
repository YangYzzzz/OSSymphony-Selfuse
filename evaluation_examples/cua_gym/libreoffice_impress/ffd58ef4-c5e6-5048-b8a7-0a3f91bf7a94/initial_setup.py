"""
Initial Setup: Build 10-slide sales pitch presentation (before objection-handling appendix).
Task ID: impress_sales_063
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_063'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title(slide, text, font_size=Pt(32), bold=True, color=None):
    """Set the title placeholder text with formatting."""
    title = slide.shapes.title
    title.text = text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=Pt(18),
                 bold=False, color=None, alignment=None):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = font_size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=Pt(16)):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = font_size
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # === Slide 1: Title Slide ===
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "NexaFlow Enterprise Platform"
    slide1.placeholders[1].text = "Transforming Business Operations Through Intelligent Automation"
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0D, 0x23, 0x3B)
    for para in slide1.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(40)
            run.font.bold = True
    for para in slide1.placeholders[1].text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
            run.font.size = Pt(22)

    # === Slide 2: Company Overview ===
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "About NexaFlow", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    add_bullet_points(slide2, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
        "Founded in 2018 by former enterprise architects from Salesforce and SAP",
        "Headquarters in San Francisco with offices in London, Singapore, and Tokyo",
        "Over 450 enterprise clients across 32 countries",
        "2024 Revenue: $187M (42% YoY growth)",
        "Named Leader in Gartner Magic Quadrant for Process Automation (2024)",
        "Team of 820+ employees including 340 engineers",
    ], Pt(16))

    # === Slide 3: The Problem ===
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "The Challenge", Pt(32), bold=True,
                 color=RGBColor(0xCC, 0x33, 0x33))
    add_bullet_points(slide3, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
        "67% of enterprises run 200+ disconnected SaaS applications",
        "Manual data entry costs the average mid-market firm $2.4M annually",
        "Integration projects take 6-18 months with traditional middleware",
        "Employee productivity loss: 28% of work hours spent on repetitive tasks",
        "Data silos lead to inconsistent reporting and delayed decision-making",
    ], Pt(16))

    # === Slide 4: Our Solution ===
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "NexaFlow: The Intelligent Automation Platform", Pt(28), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    add_bullet_points(slide4, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
        "AI-powered workflow automation across 500+ app connectors",
        "No-code visual builder with drag-and-drop process design",
        "Real-time data sync and transformation engine",
        "Pre-built industry templates for Finance, Healthcare, and Retail",
        "99.99% uptime SLA with enterprise-grade infrastructure",
    ], Pt(16))

    # === Slide 5: Key Features ===
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Platform Features", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))

    features = [
        ("Smart Connectors", "500+ pre-built integrations with auto-mapping"),
        ("AI Assistant", "Natural language workflow creation and optimization"),
        ("Analytics Dashboard", "Real-time monitoring with custom KPI tracking"),
        ("Version Control", "Full audit trail and rollback capabilities"),
    ]
    y_pos = Inches(1.8)
    for feat_name, feat_desc in features:
        add_text_box(slide5, Inches(1.0), y_pos, Inches(3), Inches(0.5),
                     feat_name, Pt(18), bold=True, color=RGBColor(0x22, 0x66, 0xAA))
        add_text_box(slide5, Inches(4.2), y_pos, Inches(7), Inches(0.5),
                     feat_desc, Pt(16))
        y_pos += Inches(1.1)

    # === Slide 6: Benefits ===
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Measurable Business Impact", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    add_bullet_points(slide6, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
        "73% reduction in manual data entry within first 90 days",
        "Average ROI of 340% within 12 months of deployment",
        "Integration deployment time reduced from months to days",
        "Employee satisfaction scores up 31% post-implementation",
        "$1.8M average annual savings per enterprise client",
    ], Pt(16))

    # === Slide 7: Pricing ===
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Pricing Plans", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))

    # Pricing table
    tbl_shape = slide7.shapes.add_table(4, 4, Inches(1), Inches(2), Inches(10), Inches(3))
    tbl = tbl_shape.table
    headers = ["Feature", "Starter", "Professional", "Enterprise"]
    data = [
        ["Connectors", "50", "200", "Unlimited"],
        ["Workflows", "25", "100", "Unlimited"],
        ["Monthly Price", "$499/mo", "$1,499/mo", "Custom"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # === Slide 8: Case Study ===
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Case Study: Meridian Healthcare Group", Pt(28), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    add_bullet_points(slide8, Inches(0.8), Inches(1.8), Inches(11), Inches(5), [
        "Challenge: 14 disconnected patient management systems across 23 facilities",
        "Solution: NexaFlow unified data pipeline with HIPAA-compliant connectors",
        "Timeline: Full deployment in 8 weeks (vs. 12-month industry average)",
        "Results: 89% reduction in data entry errors",
        "Annual savings: $3.2M in operational costs",
        "Patient wait times decreased by 34%",
    ], Pt(16))

    # === Slide 9: Team ===
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide9, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Leadership Team", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    team = [
        ("Dr. Anika Patel", "CEO & Co-Founder", "Former VP Engineering at Salesforce"),
        ("James Whitmore", "CTO & Co-Founder", "Ex-SAP Chief Architect"),
        ("Lin Wei", "VP of Product", "Previously led automation at ServiceNow"),
        ("Sofia Rodriguez", "VP of Sales", "15 years enterprise sales at Oracle"),
    ]
    y_pos = Inches(1.8)
    for name, title, bg in team:
        add_text_box(slide9, Inches(1.0), y_pos, Inches(4), Inches(0.4),
                     name, Pt(18), bold=True, color=RGBColor(0x0D, 0x23, 0x3B))
        add_text_box(slide9, Inches(1.0), y_pos + Inches(0.4), Inches(4), Inches(0.4),
                     f"{title} - {bg}", Pt(14))
        y_pos += Inches(1.1)

    # === Slide 10: Contact / Next Steps ===
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    add_text_box(slide10, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
                 "Next Steps", Pt(32), bold=True,
                 color=RGBColor(0x0D, 0x23, 0x3B))
    add_bullet_points(slide10, Inches(0.8), Inches(1.8), Inches(11), Inches(4), [
        "Schedule a personalized demo with our solutions team",
        "Start a 30-day free trial with full feature access",
        "Review our implementation playbook and onboarding timeline",
    ], Pt(18))
    add_text_box(slide10, Inches(0.8), Inches(4.5), Inches(10), Inches(1),
                 "Contact: sales@nexaflow.io | +1 (415) 555-0192",
                 Pt(16), color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.LEFT)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Launch LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
