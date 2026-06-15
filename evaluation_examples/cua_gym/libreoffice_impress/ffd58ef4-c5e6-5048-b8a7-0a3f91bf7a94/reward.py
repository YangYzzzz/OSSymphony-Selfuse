"""
Reward Script: Objection-handling appendix slides 11-14
Task ID: impress_sales_063
Domain: libreoffice_impress
Scoring:
  Component 1: Slide count increased to 14 (0.10)
  Component 2: Slide 11 - "Too Expensive" title formatting + ROI table (0.25)
  Component 3: Slide 12 - "We Already Have a Solution" title formatting + comparison table (0.25)
  Component 4: Slide 13 - "Implementation Is Too Complex" title formatting + 4-step process shapes (0.20)
  Component 5: Slide 14 - "Security Concerns" title formatting + 4 compliance shield shapes (0.20)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_063'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice edits via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_title_info(slide):
    """Extract title text and font properties from a slide.
    Searches text boxes for the objection title (bold, red, ~28pt).
    Returns (text, bold, size_pt, color_hex) or (None, None, None, None).
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if not text:
                    continue
                bold = run.font.bold
                size = run.font.size  # EMU
                try:
                    color = str(run.font.color.rgb) if run.font.color.type is not None else None
                except Exception:
                    color = None
                # Look for bold red title text
                if bold and color == "CC0000":
                    size_pt = round(size / 12700) if size else None
                    return text, bold, size_pt, color
    return None, None, None, None


def count_tables(slide):
    """Count table shapes and return list of (rows, cols) tuples."""
    tables = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            t = shape.table
            tables.append((len(t.rows), len(t.columns)))
    return tables


def count_auto_shapes(slide):
    """Count AUTO_SHAPE shapes on a slide."""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            count += 1
    return count


def get_auto_shape_texts(slide):
    """Get text from all auto shapes on a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
            full_text = ""
            for para in shape.text_frame.paragraphs:
                full_text += para.text
            texts.append(full_text.strip())
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Presentation has exactly 14 slides (0.10 points)
    # Initial has 10 slides; task adds 4 (slides 11-14)
    try:
        if num_slides >= 14:
            print(f"PASS: Component 1 — Slide count is {num_slides} (>= 14) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — Expected >= 14 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 14 slides to check the rest
    if num_slides < 14:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2a: Slide 11 - "Too Expensive" title in bold #CC0000 (0.10 points)
    try:
        slide_11 = prs.slides[10]
        title_text, title_bold, title_size, title_color = get_title_info(slide_11)

        if title_text and "too expensive" in title_text.lower() and title_color == "CC0000" and title_bold:
            print(f"PASS: Component 2a — Slide 11 title '{title_text}' bold={title_bold} size={title_size}pt color={title_color} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2a — Expected 'Too Expensive' in bold #CC0000, found: text={title_text} bold={title_bold} color={title_color}")
    except Exception as e:
        print(f"ERROR: Component 2a — {e}")

    # Component 2b: Slide 11 - ROI breakdown table (0.15 points)
    try:
        tables = count_tables(prs.slides[10])
        if len(tables) >= 1 and tables[0][0] >= 3 and tables[0][1] >= 2:
            print(f"PASS: Component 2b — ROI table found: {tables[0][0]}x{tables[0][1]} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2b — No suitable table on slide 11 (tables found: {tables})")
    except Exception as e:
        print(f"ERROR: Component 2b — {e}")

    # Component 3a: Slide 12 - "We Already Have a Solution" title in bold #CC0000 (0.10 points)
    try:
        slide_12 = prs.slides[11]
        title_text, title_bold, title_size, title_color = get_title_info(slide_12)

        if title_text and ("already have" in title_text.lower() or "have a solution" in title_text.lower()) and title_color == "CC0000" and title_bold:
            print(f"PASS: Component 3a — Slide 12 title '{title_text}' bold={title_bold} color={title_color} (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3a — Expected 'already have a solution' in bold #CC0000, found: text={title_text} bold={title_bold} color={title_color}")
    except Exception as e:
        print(f"ERROR: Component 3a — {e}")

    # Component 3b: Slide 12 - Comparison table (0.15 points)
    try:
        tables = count_tables(prs.slides[11])
        if len(tables) >= 1 and tables[0][0] >= 3 and tables[0][1] >= 3:
            print(f"PASS: Component 3b — Comparison table found: {tables[0][0]}x{tables[0][1]} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3b — No suitable comparison table on slide 12 (tables found: {tables})")
    except Exception as e:
        print(f"ERROR: Component 3b — {e}")

    # Component 4a: Slide 13 - "Implementation Is Too Complex" title in bold #CC0000 (0.08 points)
    try:
        slide_13 = prs.slides[12]
        title_text, title_bold, title_size, title_color = get_title_info(slide_13)

        if title_text and ("complex" in title_text.lower() or "implementation" in title_text.lower()) and title_color == "CC0000" and title_bold:
            print(f"PASS: Component 4a — Slide 13 title '{title_text}' bold={title_bold} color={title_color} (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 4a — Expected 'Implementation/Complex' in bold #CC0000, found: text={title_text} bold={title_bold} color={title_color}")
    except Exception as e:
        print(f"ERROR: Component 4a — {e}")

    # Component 4b: Slide 13 - 4-step process diagram shapes (0.12 points)
    try:
        slide_13 = prs.slides[12]
        auto_count = count_auto_shapes(slide_13)
        shape_texts = get_auto_shape_texts(slide_13)
        step_keywords = ["step 1", "step 2", "step 3", "step 4"]
        found_steps = sum(1 for kw in step_keywords if any(kw in t.lower() for t in shape_texts))

        if auto_count >= 4 and found_steps >= 4:
            print(f"PASS: Component 4b — 4-step process diagram: {auto_count} auto shapes, {found_steps} steps found (0.12 pts)")
            total_score += 0.12
        elif auto_count >= 4 and found_steps >= 2:
            print(f"PARTIAL: Component 4b — Found {found_steps}/4 steps in {auto_count} shapes (0.06 pts)")
            total_score += 0.06
        else:
            print(f"FAIL: Component 4b — Expected >= 4 auto shapes with step labels, found {auto_count} shapes, {found_steps} steps")
    except Exception as e:
        print(f"ERROR: Component 4b — {e}")

    # Component 5a: Slide 14 - "Security Concerns" title in bold #CC0000 (0.08 points)
    try:
        slide_14 = prs.slides[13]
        title_text, title_bold, title_size, title_color = get_title_info(slide_14)

        if title_text and "security" in title_text.lower() and title_color == "CC0000" and title_bold:
            print(f"PASS: Component 5a — Slide 14 title '{title_text}' bold={title_bold} color={title_color} (0.08 pts)")
            total_score += 0.08
        else:
            print(f"FAIL: Component 5a — Expected 'Security Concerns' in bold #CC0000, found: text={title_text} bold={title_bold} color={title_color}")
    except Exception as e:
        print(f"ERROR: Component 5a — {e}")

    # Component 5b: Slide 14 - 4 compliance badge shapes (SOC2, GDPR, HIPAA, ISO27001) (0.12 points)
    try:
        slide_14 = prs.slides[13]
        shape_texts = get_auto_shape_texts(slide_14)
        badge_keywords = ["soc2", "gdpr", "hipaa", "iso27001"]
        found_badges = sum(1 for kw in badge_keywords if any(kw in t.lower() for t in shape_texts))
        auto_count = count_auto_shapes(slide_14)

        if found_badges >= 4 and auto_count >= 4:
            print(f"PASS: Component 5b — All 4 compliance badges found in {auto_count} shapes (0.12 pts)")
            total_score += 0.12
        elif found_badges >= 2:
            print(f"PARTIAL: Component 5b — Found {found_badges}/4 badges in {auto_count} shapes (0.06 pts)")
            total_score += 0.06
        else:
            print(f"FAIL: Component 5b — Only {found_badges} compliance badges found in {auto_count} shapes")
    except Exception as e:
        print(f"ERROR: Component 5b — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
