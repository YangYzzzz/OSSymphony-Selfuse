"""
FINAL REWARD SCRIPT - SUCCESS
Task: I've got this slide that I really like, and I want to create an exact copy of it right after the original. What's the quickest way to duplicate it in LibreOffice Impress?
Generated: 2025-08-07 08:51:37
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
import hashlib
from pptx import Presentation

def get_slide_texts(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                texts.append(text)
    return texts

def get_slide_image_hashes(slide):
    hashes = []
    for shape in slide.shapes:
        if hasattr(shape, 'image') and shape.image is not None:
            blob = shape.image.blob
            h = hashlib.md5(blob).hexdigest()
            hashes.append(h)
    return hashes


def verify_slide_duplication(file_path):
    print(f"Checking slide duplication for file: {file_path}")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2 points)
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {total_score:.2f}")
        return total_score
    print(f"✓ File exists (0.2 points)")
    total_score += 0.2

    # Requirement 2: Load presentation (no direct score, prerequisite for next checks)
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print(f"REWARD: {total_score:.2f}")
        return total_score

    # Requirement 3: Slide count >= 2 (0.2 points)
    slide_count = len(prs.slides)
    if slide_count >= 2:
        print(f"✓ Slide count >= 2 ({slide_count}) (0.2 points)")
        total_score += 0.2
    else:
        print(f"✗ Slide count < 2 ({slide_count}), duplication failed")
        print(f"REWARD: {total_score:.2f}")
        return total_score

    # Compare the first two slides for duplication quality
    slide1 = prs.slides[0]
    slide2 = prs.slides[1]

    # Requirement 4: Layout match (0.1 points)
    layout1 = getattr(slide1.slide_layout, 'name', None)
    layout2 = getattr(slide2.slide_layout, 'name', None)
    if layout1 and layout1 == layout2:
        print(f"✓ Slide layouts match ('{layout1}') (0.1 points)")
        total_score += 0.1
    else:
        print(f"✗ Slide layouts differ ('{layout1}' vs '{layout2}')")

    # Requirement 5: Text content match (0.3 points)
    texts1 = get_slide_texts(slide1)
    texts2 = get_slide_texts(slide2)
    matched_texts = sum(1 for t in texts1 if t in texts2)
    if texts1:
        text_score = 0.3 * (matched_texts / len(texts1))
    else:
        # No text on slide considered full match
        text_score = 0.3
    total_score += text_score
    print(f"Found {len(texts1)} texts in slide1, {matched_texts} matched (Text score: {text_score:.2f})")

    # Requirement 6: Image content match (0.2 points)
    hashes1 = get_slide_image_hashes(slide1)
    hashes2 = get_slide_image_hashes(slide2)
    matched_imgs = sum(1 for h in hashes1 if h in hashes2)
    if hashes1:
        img_score = 0.2 * (matched_imgs / len(hashes1))
    else:
        # No images considered full match
        img_score = 0.2
    total_score += img_score
    print(f"Found {len(hashes1)} images in slide1, {matched_imgs} matched (Image score: {img_score:.2f})")

    # Final score calculation
    final_score = min(total_score, max_score)
    print(f"Total score: {final_score:.2f}/{max_score}")
    print(f"REWARD: {final_score:.2f}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/ive_got_this_slide_that_i_really_like_and_i_want_to_create_an_exact_copy_of_it_right_after_the_origi.pptx'
    verify_slide_duplication(file_path)
