"""
Reward Script: Enable auto-fit text on overflowing text boxes
Task ID: impress_fix_015
Domain: libreoffice_impress
Scoring: 0.25 per slide (slides 2, 4, 7, 9) where the overflow text box
         has normAutofit enabled (shrink text on overflow).
"""

import os
from pptx import Presentation
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_015'

# Target slides (1-indexed) and the shape_id that should have auto-fit changed
TARGET_SLIDES = [2, 4, 7, 9]
POINTS_PER_SLIDE = 0.25


def find_overflow_textbox(slide):
    """
    Find the content text box on a target slide.
    The overflow text box is the one that was originally noAutofit in initial_env.
    We identify it as the non-title text shape with the most text content.
    """
    best_shape = None
    best_len = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip title placeholders (usually short)
        text = shape.text or ""
        # We want the content body text box, not the title
        # The target shapes have substantial text (multi-line content)
        if len(text) > best_len and len(text) > 50:
            best_len = len(text)
            best_shape = shape
    return best_shape


def check_autofit(shape):
    """
    Check if a shape has normAutofit (shrink text on overflow) enabled.
    Returns: (bool, str) - (is_normAutofit, actual_type)
    """
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return False, 'no_bodyPr'

    normAuto = bodyPr.find(qn('a:normAutofit'))
    spAuto = bodyPr.find(qn('a:spAutoFit'))
    noAuto = bodyPr.find(qn('a:noAutofit'))

    if normAuto is not None:
        return True, 'normAutofit'
    elif spAuto is not None:
        return False, 'spAutoFit'
    elif noAuto is not None:
        return False, 'noAutofit'
    else:
        return False, 'none/inherit'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    if num_slides < 10:
        print(f"FAIL: Expected at least 10 slides, found {num_slides}")
        print("REWARD: 0.0")
        return 0.0

    for slide_num in TARGET_SLIDES:
        slide_idx = slide_num - 1
        slide = prs.slides[slide_idx]

        # Component: Slide N overflow text box has normAutofit enabled
        try:
            target_shape = find_overflow_textbox(slide)
            if target_shape is None:
                print(f"FAIL: Slide {slide_num} - no content text box found")
                continue

            is_autofit, actual_type = check_autofit(target_shape)
            text_preview = (target_shape.text or "")[:40]

            if is_autofit:
                print(f"PASS: Slide {slide_num} - auto-fit enabled (normAutofit) on '{text_preview}...' ({POINTS_PER_SLIDE} pts)")
                total_score += POINTS_PER_SLIDE
            else:
                print(f"FAIL: Slide {slide_num} - expected normAutofit, found {actual_type} on '{text_preview}...'")
        except Exception as e:
            print(f"ERROR: Slide {slide_num} - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
