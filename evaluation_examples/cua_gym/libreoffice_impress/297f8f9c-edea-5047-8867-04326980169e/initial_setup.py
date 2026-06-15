"""
Initial Setup: Create a quarterly report presentation with text overflow issues
Task ID: impress_fix_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_textbox(tf, text, font_size=Pt(14), font_name="Arial",
                        bold=False, color=None, alignment=None):
    """Helper to add formatted text to a text frame."""
    p = tf.paragraphs[0]
    if alignment:
        p.alignment = alignment
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = font_name
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_paragraph(tf, text, font_size=Pt(14), font_name="Arial",
                  bold=False, color=None, level=0):
    """Add a new paragraph to a text frame."""
    p = tf.add_paragraph()
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.name = font_name
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    dark_blue = RGBColor(0x1B, 0x3A, 0x5C)
    medium_blue = RGBColor(0x2E, 0x75, 0xB6)
    dark_gray = RGBColor(0x33, 0x33, 0x33)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    light_gray = RGBColor(0x66, 0x66, 0x66)

    # ========== SLIDE 1: Title Slide ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = dark_blue

    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Quarterly Business Review"
    run.font.size = Pt(44)
    run.font.name = "Arial"
    run.font.bold = True
    run.font.color.rgb = white

    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Q4 2025 - TechVenture Solutions Inc."
    run2.font.size = Pt(24)
    run2.font.name = "Arial"
    run2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)

    txBox3 = slide1.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(0.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Prepared by: Strategic Planning Division | December 2025"
    run3.font.size = Pt(14)
    run3.font.name = "Arial"
    run3.font.color.rgb = RGBColor(0x99, 0xAA, 0xBB)

    # ========== SLIDE 2: Executive Summary (OVERFLOW) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t2 = title2.text_frame
    tf_t2.word_wrap = True
    p_t2 = tf_t2.paragraphs[0]
    run_t2 = p_t2.add_run()
    run_t2.text = "Executive Summary"
    run_t2.font.size = Pt(32)
    run_t2.font.bold = True
    run_t2.font.color.rgb = dark_blue

    # OVERFLOW text box - small box with lots of text, auto-fit DISABLED
    overflow2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.0))
    tf_o2 = overflow2.text_frame
    tf_o2.word_wrap = True
    tf_o2.auto_size = MSO_AUTO_SIZE.NONE  # Explicitly disable auto-fit
    p = tf_o2.paragraphs[0]
    run = p.add_run()
    run.text = "Revenue Performance Overview"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = medium_blue

    overflow_texts = [
        "Total revenue for Q4 2025 reached $47.8 million, representing a 23% year-over-year increase compared to Q4 2024. This exceeds our internal target of $44.5 million by 7.4%.",
        "The Enterprise Solutions segment contributed $28.3 million (59% of total), driven by three major contract renewals with Fortune 500 clients: Meridian Healthcare ($4.2M), Atlas Financial ($3.8M), and Pinnacle Manufacturing ($2.9M).",
        "Our SaaS platform subscriptions grew to 12,847 active accounts, up from 9,230 in Q3. Monthly recurring revenue (MRR) increased to $3.92 million, putting us on track for $47M ARR.",
        "International markets showed particularly strong performance, with APAC revenue growing 41% and EMEA growing 18%. The newly opened Singapore office contributed $2.1M in its first full quarter.",
        "Gross margins improved to 72.3% from 68.9% in Q3, primarily due to infrastructure optimization and improved vendor negotiations. Operating expenses remained flat at $31.2 million despite headcount growth of 8%."
    ]
    for txt in overflow_texts:
        add_paragraph(tf_o2, txt, font_size=Pt(13), color=dark_gray)

    # Right side text box (normal, not overflow)
    right2 = slide2.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(5))
    tf_r2 = right2.text_frame
    tf_r2.word_wrap = True
    p_r2 = tf_r2.paragraphs[0]
    run_r2 = p_r2.add_run()
    run_r2.text = "Key Metrics at a Glance"
    run_r2.font.size = Pt(18)
    run_r2.font.bold = True
    run_r2.font.color.rgb = medium_blue

    metrics = [
        "Revenue: $47.8M (+23% YoY)",
        "Gross Margin: 72.3%",
        "Active SaaS Accounts: 12,847",
        "Customer Retention: 94.2%",
        "Net Promoter Score: 67",
        "Employee Count: 423"
    ]
    for m in metrics:
        add_paragraph(tf_r2, m, font_size=Pt(14), color=dark_gray)

    # ========== SLIDE 3: Financial Highlights ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t3 = title3.text_frame
    tf_t3.word_wrap = True
    p_t3 = tf_t3.paragraphs[0]
    run_t3 = p_t3.add_run()
    run_t3.text = "Financial Highlights - Q4 2025"
    run_t3.font.size = Pt(32)
    run_t3.font.bold = True
    run_t3.font.color.rgb = dark_blue

    # Table with financial data
    table_shape = slide3.shapes.add_table(7, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5))
    table = table_shape.table
    headers = ["Category", "Q3 2025", "Q4 2025", "Change"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    fin_data = [
        ["Total Revenue", "$38.9M", "$47.8M", "+22.9%"],
        ["Cost of Revenue", "$12.1M", "$13.2M", "+9.1%"],
        ["Gross Profit", "$26.8M", "$34.6M", "+29.1%"],
        ["Operating Expenses", "$30.8M", "$31.2M", "+1.3%"],
        ["EBITDA", "$5.2M", "$8.4M", "+61.5%"],
        ["Net Income", "$3.1M", "$5.7M", "+83.9%"],
    ]
    for r, row in enumerate(fin_data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # ========== SLIDE 4: Product Performance (OVERFLOW) ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t4 = title4.text_frame
    tf_t4.word_wrap = True
    p_t4 = tf_t4.paragraphs[0]
    run_t4 = p_t4.add_run()
    run_t4.text = "Product Line Performance"
    run_t4.font.size = Pt(32)
    run_t4.font.bold = True
    run_t4.font.color.rgb = dark_blue

    # OVERFLOW text box
    overflow4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(2.2))
    tf_o4 = overflow4.text_frame
    tf_o4.word_wrap = True
    tf_o4.auto_size = MSO_AUTO_SIZE.NONE  # Explicitly disable auto-fit

    p = tf_o4.paragraphs[0]
    run = p.add_run()
    run.text = "Detailed Product Analysis"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = medium_blue

    product_texts = [
        "CloudSync Enterprise (Revenue: $18.7M | Growth: +31%): Our flagship cloud integration platform continued its strong trajectory. The release of version 4.2 in October introduced real-time data synchronization across hybrid cloud environments, which was the primary driver for three new enterprise deals.",
        "DataVault Analytics (Revenue: $12.4M | Growth: +19%): The business intelligence suite saw steady adoption among mid-market clients. The new predictive analytics module launched in November has already generated $1.8M in pipeline, with 24 active proof-of-concept deployments.",
        "SecureGate Identity (Revenue: $9.8M | Growth: +28%): Zero-trust identity management became our fastest-growing product line. The partnership with CyberShield Technologies expanded our threat detection capabilities, resulting in a 340% increase in inbound demo requests.",
        "WorkFlow Automation (Revenue: $6.9M | Growth: +12%): Process automation tools showed moderate growth. The integration with major ERP platforms (SAP, Oracle, Microsoft Dynamics) completed in Q4, opening access to an estimated $2.3 billion addressable market segment."
    ]
    for txt in product_texts:
        add_paragraph(tf_o4, txt, font_size=Pt(13), color=dark_gray)

    # ========== SLIDE 5: Customer Analysis ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t5 = title5.text_frame
    tf_t5.word_wrap = True
    p_t5 = tf_t5.paragraphs[0]
    run_t5 = p_t5.add_run()
    run_t5.text = "Customer Segmentation Analysis"
    run_t5.font.size = Pt(32)
    run_t5.font.bold = True
    run_t5.font.color.rgb = dark_blue

    left5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5))
    tf_l5 = left5.text_frame
    tf_l5.word_wrap = True
    p_l5 = tf_l5.paragraphs[0]
    run_l5 = p_l5.add_run()
    run_l5.text = "Enterprise (500+ employees)"
    run_l5.font.size = Pt(16)
    run_l5.font.bold = True
    run_l5.font.color.rgb = medium_blue
    ent_items = [
        "Active accounts: 187 (+14 new in Q4)",
        "Average contract value: $152K/year",
        "Retention rate: 97.3%",
        "Upsell rate: 34% purchased additional modules"
    ]
    for item in ent_items:
        add_paragraph(tf_l5, item, font_size=Pt(13), color=dark_gray)

    add_paragraph(tf_l5, "Mid-Market (50-499 employees)", font_size=Pt(16),
                  bold=True, color=medium_blue)
    mid_items = [
        "Active accounts: 1,243 (+89 new in Q4)",
        "Average contract value: $38K/year",
        "Retention rate: 91.8%"
    ]
    for item in mid_items:
        add_paragraph(tf_l5, item, font_size=Pt(13), color=dark_gray)

    right5 = slide5.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5))
    tf_r5 = right5.text_frame
    tf_r5.word_wrap = True
    p_r5 = tf_r5.paragraphs[0]
    run_r5 = p_r5.add_run()
    run_r5.text = "SMB (Under 50 employees)"
    run_r5.font.size = Pt(16)
    run_r5.font.bold = True
    run_r5.font.color.rgb = medium_blue
    smb_items = [
        "Active accounts: 11,417 (+1,206 new in Q4)",
        "Average contract value: $8.4K/year",
        "Retention rate: 86.4%",
        "Self-service onboarding: 78% of new signups"
    ]
    for item in smb_items:
        add_paragraph(tf_r5, item, font_size=Pt(13), color=dark_gray)

    # ========== SLIDE 6: Team & Hiring ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t6 = title6.text_frame
    tf_t6.word_wrap = True
    p_t6 = tf_t6.paragraphs[0]
    run_t6 = p_t6.add_run()
    run_t6.text = "People & Culture Update"
    run_t6.font.size = Pt(32)
    run_t6.font.bold = True
    run_t6.font.color.rgb = dark_blue

    content6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_c6 = content6.text_frame
    tf_c6.word_wrap = True
    p_c6 = tf_c6.paragraphs[0]
    run_c6 = p_c6.add_run()
    run_c6.text = "Headcount grew from 392 to 423 employees during Q4, with key hires in Engineering (12), Sales (8), Customer Success (6), and Product (5)."
    run_c6.font.size = Pt(14)
    run_c6.font.color.rgb = dark_gray

    team_updates = [
        "Employee satisfaction score: 4.3/5.0 (up from 4.1 in Q3)",
        "Voluntary attrition rate: 8.2% annualized (industry average: 13.5%)",
        "Diversity hiring: 47% of new hires from underrepresented groups",
        "Remote/hybrid workforce: 62% hybrid, 28% fully remote, 10% office-based",
        "Training investment: $1,850 per employee in Q4 professional development"
    ]
    for item in team_updates:
        add_paragraph(tf_c6, item, font_size=Pt(14), color=dark_gray)

    # ========== SLIDE 7: Strategic Initiatives (OVERFLOW) ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t7 = title7.text_frame
    tf_t7.word_wrap = True
    p_t7 = tf_t7.paragraphs[0]
    run_t7 = p_t7.add_run()
    run_t7.text = "Strategic Initiatives & Roadmap"
    run_t7.font.size = Pt(32)
    run_t7.font.bold = True
    run_t7.font.color.rgb = dark_blue

    # OVERFLOW text box
    overflow7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(2.0))
    tf_o7 = overflow7.text_frame
    tf_o7.word_wrap = True
    tf_o7.auto_size = MSO_AUTO_SIZE.NONE  # Explicitly disable auto-fit

    p = tf_o7.paragraphs[0]
    run = p.add_run()
    run.text = "Q1 2026 Strategic Priorities"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = medium_blue

    strategy_texts = [
        "AI Integration Initiative (Budget: $4.5M): Deploy machine learning capabilities across all product lines. Phase 1 focuses on predictive customer churn modeling and intelligent ticket routing for CloudSync. Expected to reduce customer support costs by 25% while improving response accuracy.",
        "Global Expansion - Latin America (Budget: $3.2M): Establish regional presence in Sao Paulo and Mexico City. Initial target is 50 enterprise accounts by Q2 2026. Local language support for Portuguese and Spanish will be completed by end of January.",
        "Platform Consolidation (Budget: $2.8M): Merge DataVault and WorkFlow Automation into a unified analytics and automation platform. This reduces infrastructure costs by an estimated $1.4M annually and simplifies the customer upgrade path.",
        "Security Certification Drive: Achieve SOC 2 Type II, ISO 27001, and FedRAMP Moderate certifications by Q3 2026. This unlocks an estimated $8.5M in government and regulated industry pipeline that currently requires these compliance standards.",
        "Partner Ecosystem Growth: Expand systems integrator partnerships from 12 to 25. Launch certified partner training program with tiered benefits. Target $6M in partner-influenced revenue for 2026."
    ]
    for txt in strategy_texts:
        add_paragraph(tf_o7, txt, font_size=Pt(13), color=dark_gray)

    # ========== SLIDE 8: Risk Assessment ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t8 = title8.text_frame
    tf_t8.word_wrap = True
    p_t8 = tf_t8.paragraphs[0]
    run_t8 = p_t8.add_run()
    run_t8.text = "Risk Assessment & Mitigation"
    run_t8.font.size = Pt(32)
    run_t8.font.bold = True
    run_t8.font.color.rgb = dark_blue

    content8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
    tf_c8 = content8.text_frame
    tf_c8.word_wrap = True
    p_c8 = tf_c8.paragraphs[0]
    run_c8 = p_c8.add_run()
    run_c8.text = "Key risks identified for the upcoming quarter include market competition intensification, particularly from CloudNova's recent Series D funding of $200M, and potential macroeconomic headwinds affecting enterprise IT budgets."
    run_c8.font.size = Pt(14)
    run_c8.font.color.rgb = dark_gray

    risk_items = [
        "Competitive pressure: CloudNova and DataBridge both launched competing products in Q4",
        "Talent retention: Engineering salaries continue to rise; 3 senior engineers received competing offers",
        "Currency exposure: 28% of revenue in non-USD currencies; hedging strategy under review",
        "Regulatory: EU AI Act compliance requirements being assessed for ML features"
    ]
    for item in risk_items:
        add_paragraph(tf_c8, item, font_size=Pt(14), color=dark_gray)

    # ========== SLIDE 9: Q1 2026 Outlook (OVERFLOW) ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf_t9 = title9.text_frame
    tf_t9.word_wrap = True
    p_t9 = tf_t9.paragraphs[0]
    run_t9 = p_t9.add_run()
    run_t9.text = "Q1 2026 Outlook & Targets"
    run_t9.font.size = Pt(32)
    run_t9.font.bold = True
    run_t9.font.color.rgb = dark_blue

    # OVERFLOW text box
    overflow9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(2.0))
    tf_o9 = overflow9.text_frame
    tf_o9.word_wrap = True
    tf_o9.auto_size = MSO_AUTO_SIZE.NONE  # Explicitly disable auto-fit

    p = tf_o9.paragraphs[0]
    run = p.add_run()
    run.text = "Financial Targets & Operational Goals"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = medium_blue

    outlook_texts = [
        "Revenue target: $51.2M (+7.1% QoQ). This assumes continued enterprise momentum with 8 renewals worth $14.2M in aggregate, plus new logo acquisition targets of 15 enterprise and 120 mid-market accounts. Seasonality adjustment applied for Q1 historical softness.",
        "Product launches: CloudSync 5.0 beta release scheduled for February 15th with general availability on March 28th. Key features include native Kubernetes orchestration, multi-region failover, and real-time collaborative dashboards built on our new WebSocket architecture.",
        "Operational efficiency: Target operating margin improvement to 15.2% from current 13.8%. Key levers include automation of Tier 1 support (projected 40% deflection rate), renegotiated AWS reserved instance pricing (-18% on compute), and consolidation of three legacy monitoring tools.",
        "Customer success: Reduce average onboarding time from 14 days to 8 days through the new guided setup wizard. Increase product adoption score from 6.2 to 7.0 on our 10-point engagement index. Launch customer advisory board with 12 strategic accounts.",
        "Market expansion: Complete localization for Japanese and Korean markets. Attend 4 major industry conferences (AWS re:Invent, Gartner Symposium, RSA Conference, SaaStr Annual). Launch referral program targeting 200 qualified leads per month."
    ]
    for txt in outlook_texts:
        add_paragraph(tf_o9, txt, font_size=Pt(13), color=dark_gray)

    # ========== SLIDE 10: Closing ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    fill10 = slide10.background.fill
    fill10.solid()
    fill10.fore_color.rgb = dark_blue

    txBox10 = slide10.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf10 = txBox10.text_frame
    tf10.word_wrap = True
    p10 = tf10.paragraphs[0]
    p10.alignment = PP_ALIGN.CENTER
    run10 = p10.add_run()
    run10.text = "Thank You"
    run10.font.size = Pt(44)
    run10.font.bold = True
    run10.font.color.rgb = white

    txBox10b = slide10.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
    tf10b = txBox10b.text_frame
    tf10b.word_wrap = True
    p10b = tf10b.paragraphs[0]
    p10b.alignment = PP_ALIGN.CENTER
    run10b = p10b.add_run()
    run10b.text = "Questions? Contact: strategy@techventuresolutions.com"
    run10b.font.size = Pt(18)
    run10b.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
