"""
Initial Setup: Chemistry Lecture presentation with 10 slides, no transitions
Task ID: impress_teach_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.paragraphs[0].text = point
        else:
            p = tf.add_paragraph()
            p.text = point
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Introduction to Organic Chemistry",
        "CHEM 201 — Dr. Elena Vasquez\nFall 2025 Semester"
    )

    # Slide 2: Course Overview
    add_content_slide(prs, "Course Overview", [
        "Fundamentals of carbon-based molecular structures",
        "Reaction mechanisms and stereochemistry",
        "Spectroscopy and analytical techniques",
        "Laboratory safety and experimental methods",
        "Prerequisites: CHEM 101 and CHEM 102"
    ])

    # Slide 3: Atomic Bonding Review
    add_content_slide(prs, "Atomic Bonding Review", [
        "Covalent bonds: sharing of electron pairs",
        "Electronegativity differences determine bond polarity",
        "Sigma (σ) bonds: head-on orbital overlap",
        "Pi (π) bonds: lateral orbital overlap",
        "Hybridization: sp³, sp², sp orbitals"
    ])

    # Slide 4: Functional Groups
    add_content_slide(prs, "Common Functional Groups", [
        "Hydroxyl (-OH): Alcohols and phenols",
        "Carbonyl (C=O): Aldehydes and ketones",
        "Carboxyl (-COOH): Carboxylic acids",
        "Amino (-NH₂): Amines and amino acids",
        "Ester (-COO-): Fats, oils, and polyesters"
    ])

    # Slide 5: Nomenclature Rules
    add_content_slide(prs, "IUPAC Nomenclature", [
        "Identify the longest carbon chain (parent chain)",
        "Number carbons from the end nearest a substituent",
        "Name substituents with appropriate prefixes",
        "Use di-, tri-, tetra- for multiple identical groups",
        "Alphabetize substituent names in the final name"
    ])

    # Slide 6: Reaction Mechanisms
    add_content_slide(prs, "Reaction Mechanisms", [
        "Nucleophilic substitution: SN1 and SN2 pathways",
        "Elimination reactions: E1 and E2 mechanisms",
        "Electrophilic addition to alkenes and alkynes",
        "Free radical chain reactions: initiation, propagation, termination",
        "Curved arrow notation for electron movement"
    ])

    # Slide 7: Stereochemistry
    add_content_slide(prs, "Stereochemistry Essentials", [
        "Chirality and asymmetric carbon centers",
        "R/S configuration using Cahn-Ingold-Prelog rules",
        "Enantiomers: non-superimposable mirror images",
        "Diastereomers: stereoisomers that are not enantiomers",
        "Optical activity and polarimetry measurements"
    ])

    # Slide 8: Spectroscopy Methods
    add_content_slide(prs, "Spectroscopy Techniques", [
        "IR Spectroscopy: functional group identification",
        "¹H NMR: proton chemical shifts and splitting",
        "¹³C NMR: carbon framework analysis",
        "Mass Spectrometry: molecular weight and fragmentation",
        "UV-Vis: conjugation and chromophore detection"
    ])

    # Slide 9: Laboratory Component
    add_content_slide(prs, "Laboratory Schedule", [
        "Week 1-3: Extraction and purification techniques",
        "Week 4-6: Synthesis of aspirin and ester derivatives",
        "Week 7-9: Grignard reaction and aldol condensation",
        "Week 10-12: Multi-step synthesis project",
        "Lab reports due one week after each experiment"
    ])

    # Slide 10: Assessment & Grading
    add_content_slide(prs, "Assessment & Grading", [
        "Midterm Exam: 25% (covers Chapters 1-6)",
        "Final Exam: 30% (comprehensive)",
        "Laboratory Reports: 20%",
        "Problem Sets: 15% (weekly submissions)",
        "Class Participation: 10%"
    ])

    # No transitions set — this is the initial state
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
