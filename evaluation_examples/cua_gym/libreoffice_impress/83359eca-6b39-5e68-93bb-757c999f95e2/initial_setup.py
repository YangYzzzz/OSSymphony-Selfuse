"""
Initial Setup: Geography Lecture presentation with 7 slides.
Slide 3 has an image on the left and a text box on the right, NO animations.
Task ID: impress_teach_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
IMG_PATH = f'{WORKDIR}/geography_map.png'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_geography_image():
    """Create a simple geography-themed placeholder image."""
    img = PILImage.new('RGB', (640, 480), color=(34, 139, 34))  # forest green
    # Draw some simple shapes to make it look like a map
    from PIL import ImageDraw, ImageFont
    draw = PILImage.core  # just use basic drawing
    d = ImageDraw.Draw(img)
    # Draw some land masses / shapes
    d.rectangle([50, 50, 300, 250], fill=(210, 180, 140))  # tan landmass
    d.rectangle([350, 100, 580, 350], fill=(210, 180, 140))  # another landmass
    d.ellipse([100, 280, 250, 420], fill=(0, 105, 148))  # lake
    d.rectangle([0, 0, 640, 30], fill=(70, 70, 70))
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 18)
    except Exception:
        font = ImageFont.load_default()
    d.text((150, 5), "World Physical Geography Map", fill=(255, 255, 255), font=font)
    img.save(IMG_PATH)
    print(f'Image created: {IMG_PATH}')


def add_text_to_shape(shape, text, font_size=18, bold=False, color=None):
    """Helper to set text on a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def create_initial():
    create_geography_image()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Geography Lecture Series"
    slide1.placeholders[1].text = "Physical Geography & Earth Systems\nProfessor Elena Vasquez\nSpring 2026"

    # ===== Slide 2: Introduction =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(title2, "Course Overview", font_size=32, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))

    body2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf2 = body2.text_frame
    tf2.word_wrap = True
    topics = [
        "This lecture series covers the fundamental concepts of physical geography,",
        "including plate tectonics, climate systems, biomes, and geomorphology.",
        "",
        "Key Topics:",
        "  - Plate Tectonics and Continental Drift",
        "  - Atmospheric Circulation and Climate Zones",
        "  - River Systems and Erosion Patterns",
        "  - Volcanic Activity and Mountain Formation",
        "  - Ocean Currents and Their Climate Effects",
    ]
    tf2.paragraphs[0].text = topics[0]
    for line in topics[1:]:
        p = tf2.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(16)
    for run in tf2.paragraphs[0].runs:
        run.font.size = Pt(16)

    # ===== Slide 3: KEY SLIDE - Image + Text, NO animations =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    add_text_to_shape(title3, "Tectonic Plate Boundaries", font_size=30, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))

    # Image on the left
    pic3 = slide3.shapes.add_picture(IMG_PATH, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5))

    # Text box on the right
    txt3 = slide3.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6), Inches(5))
    tf3 = txt3.text_frame
    tf3.word_wrap = True
    content3 = [
        "Convergent Boundaries",
        "When two plates collide, the denser oceanic plate subducts beneath the continental plate, forming deep ocean trenches and volcanic mountain ranges.",
        "",
        "Divergent Boundaries",
        "At mid-ocean ridges, magma rises to fill the gap as plates move apart, creating new oceanic crust and widening the sea floor.",
        "",
        "Transform Boundaries",
        "Plates slide horizontally past each other along transform faults. The San Andreas Fault is the most well-known example.",
    ]
    tf3.paragraphs[0].text = content3[0]
    tf3.paragraphs[0].runs[0].font.size = Pt(16)
    tf3.paragraphs[0].runs[0].font.bold = True
    for line in content3[1:]:
        p = tf3.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(14)
            if line in ("Divergent Boundaries", "Transform Boundaries"):
                run.font.bold = True
                run.font.size = Pt(16)

    # ===== Slide 4: Climate Zones =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(title4, "Global Climate Zones", font_size=30, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))
    body4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf4 = body4.text_frame
    tf4.word_wrap = True
    zones = [
        ("Tropical (0-23.5 N/S)", "High temperatures year-round with abundant rainfall. Supports dense rainforests and diverse ecosystems."),
        ("Temperate (23.5-66.5 N/S)", "Four distinct seasons with moderate precipitation. Includes deciduous forests and grasslands."),
        ("Polar (66.5-90 N/S)", "Extremely cold temperatures with minimal precipitation. Permafrost covers most of the landscape."),
        ("Arid/Semi-Arid", "Low rainfall and high evaporation rates. Found in subtropical zones and continental interiors."),
    ]
    tf4.paragraphs[0].text = f"{zones[0][0]}: {zones[0][1]}"
    tf4.paragraphs[0].runs[0].font.size = Pt(14)
    for name, desc in zones[1:]:
        p = tf4.add_paragraph()
        p.text = f"{name}: {desc}"
        for run in p.runs:
            run.font.size = Pt(14)

    # ===== Slide 5: River Systems =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(title5, "Major River Systems", font_size=30, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))

    # Table of rivers
    table_shape = slide5.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11), Inches(4))
    table = table_shape.table
    headers = ["River", "Length (km)", "Continent", "Discharge (m3/s)"]
    rivers = [
        ["Amazon", "6,992", "South America", "209,000"],
        ["Nile", "6,650", "Africa", "2,830"],
        ["Yangtze", "6,300", "Asia", "30,166"],
        ["Mississippi", "6,275", "North America", "16,800"],
        ["Yenisei", "5,539", "Asia", "19,600"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(rivers, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # ===== Slide 6: Volcanic Activity =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(title6, "Volcanic Activity & the Ring of Fire", font_size=30, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))
    body6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf6 = body6.text_frame
    tf6.word_wrap = True
    volcano_text = [
        "The Ring of Fire is a 40,000 km horseshoe-shaped zone of intense seismic and volcanic activity.",
        "It stretches from New Zealand, along the eastern edge of Asia, north across the Aleutian Islands,",
        "and south along the west coast of the Americas.",
        "",
        "Key Statistics:",
        "  - Contains 75% of the world's active volcanoes",
        "  - Responsible for 90% of earthquakes worldwide",
        "  - Home to 452 volcanoes across 15 countries",
    ]
    tf6.paragraphs[0].text = volcano_text[0]
    tf6.paragraphs[0].runs[0].font.size = Pt(14)
    for line in volcano_text[1:]:
        p = tf6.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(14)

    # ===== Slide 7: Summary =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    add_text_to_shape(title7, "Summary & Next Lecture", font_size=30, bold=True,
                      color=RGBColor(0x1B, 0x4F, 0x72))
    body7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
    tf7 = body7.text_frame
    tf7.word_wrap = True
    summary = [
        "Today we covered the fundamental building blocks of physical geography:",
        "  - Plate tectonics and boundary types",
        "  - Global climate classification systems",
        "  - Major river systems and their significance",
        "  - Volcanic activity and the Ring of Fire",
        "",
        "Next Lecture: Ocean Currents, Thermohaline Circulation, and Marine Ecosystems",
        "Reading: Chapters 8-10 of 'Physical Geography: A Landscape Appreciation'",
    ]
    tf7.paragraphs[0].text = summary[0]
    tf7.paragraphs[0].runs[0].font.size = Pt(16)
    for line in summary[1:]:
        p = tf7.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
