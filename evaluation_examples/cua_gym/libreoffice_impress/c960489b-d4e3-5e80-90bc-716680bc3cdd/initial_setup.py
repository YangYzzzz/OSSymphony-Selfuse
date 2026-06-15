"""
Initial Setup: Corporate presentation with solid blue master background
Task ID: impress_ma_029
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_029'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'
BG_IMAGE = f'{WORKDIR}/Desktop/corporate_bg.jpg'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_background_image():
    """Create a corporate background image (1920x1080) on the desktop."""
    os.makedirs(f'{WORKDIR}/Desktop', exist_ok=True)
    # Create a gradient corporate background image
    img = Image.new('RGB', (1920, 1080))
    pixels = img.load()
    for y in range(1080):
        for x in range(1920):
            # Dark blue to teal gradient with subtle pattern
            r = int(0 + (x / 1920) * 30 + (y / 1080) * 20)
            g = int(40 + (x / 1920) * 60 + (y / 1080) * 40)
            b = int(80 + (x / 1920) * 40 + (y / 1080) * 60)
            pixels[x, y] = (min(r, 255), min(g, 255), min(b, 255))
    img.save(BG_IMAGE, 'JPEG', quality=90)
    print(f'Background image created: {BG_IMAGE}')


def create_initial():
    prs = Presentation()
    # Standard 16:9 slide size
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Set master slide background to solid blue (#003366)
    master = prs.slide_masters[0]
    fill = master.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)

    # Slide content data
    slide_data = [
        {
            'layout': 0,  # Title Slide
            'title': 'Q4 2025 Corporate Strategy Review',
            'subtitle': 'Meridian Technologies Inc.\nPrepared by the Strategic Planning Division',
        },
        {
            'layout': 1,  # Title + Content
            'title': 'Executive Summary',
            'content': '• Revenue grew 18.3% year-over-year to $247.5M\n• Operating margin expanded to 24.7% from 21.2%\n• Three new enterprise clients onboarded in APAC region\n• Cloud migration initiative completed ahead of schedule\n• Employee satisfaction score reached 87/100',
        },
        {
            'layout': 1,
            'title': 'Financial Performance Overview',
            'content': '• Total Revenue: $247.5M (Target: $235M)\n• Gross Profit: $168.3M (68% margin)\n• EBITDA: $72.1M (up 22% YoY)\n• Free Cash Flow: $54.8M\n• R&D Investment: $38.2M (15.4% of revenue)',
        },
        {
            'layout': 1,
            'title': 'Market Expansion Strategy',
            'content': '• Southeast Asia office launched in Singapore\n• Partnership with Kyoto Digital Solutions finalized\n• Latin America pilot program showing 34% conversion\n• European regulatory compliance achieved (GDPR+)\n• Target: 15 new markets by end of 2026',
        },
        {
            'layout': 1,
            'title': 'Product Development Roadmap',
            'content': '• Platform v4.0 release scheduled for March 2026\n• AI-powered analytics module in beta testing\n• Mobile SDK adoption up 156% quarter-over-quarter\n• API gateway throughput improved to 50K req/sec\n• Customer-requested features backlog reduced by 40%',
        },
        {
            'layout': 1,
            'title': 'Team & Talent Metrics',
            'content': '• Headcount: 1,247 employees across 8 offices\n• Engineering team grew by 89 new hires\n• Voluntary attrition rate: 6.2% (industry avg: 13.5%)\n• Internal promotion rate: 28%\n• Diversity index improved to 0.74',
        },
        {
            'layout': 1,
            'title': 'Customer Success Highlights',
            'content': '• Net Promoter Score: 72 (up from 64)\n• Enterprise client retention: 96.8%\n• Average deal size increased to $385K\n• Support ticket resolution time: 4.2 hours\n• Onboarding time reduced by 35%',
        },
        {
            'layout': 1,
            'title': 'Risk Assessment & Mitigation',
            'content': '• Supply chain diversification completed\n• Cybersecurity audit passed with zero critical findings\n• Business continuity plan tested successfully\n• Currency hedging strategy covering 80% of FX exposure\n• Compliance team expanded for regulatory changes',
        },
        {
            'layout': 1,
            'title': 'Sustainability & ESG Initiatives',
            'content': '• Carbon footprint reduced by 23% vs baseline\n• 100% renewable energy for data centers achieved\n• Community investment fund disbursed $2.1M\n• Supply chain sustainability scoring implemented\n• ESG rating upgraded to AA by MSCI',
        },
        {
            'layout': 1,
            'title': 'Next Steps & Key Decisions',
            'content': '• Board approval needed for Series C acquisition\n• Q1 2026 budget allocation by December 15\n• Technology stack evaluation committee formation\n• Strategic partnership review with Nakamura Group\n• Annual planning offsite scheduled for January 8-10',
        },
    ]

    for i, sd in enumerate(slide_data):
        layout_idx = sd['layout']
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = sd['title']
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(32) if layout_idx == 0 else Pt(28)
                run.font.bold = True

        # Set content
        if layout_idx == 0 and 'subtitle' in sd:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.text = sd['subtitle']
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                        run.font.size = Pt(18)
        elif 'content' in sd:
            if len(slide.placeholders) > 1:
                ph = slide.placeholders[1]
                ph.text = sd['content']
                for para in ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                        run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


# Create background image first, then presentation
create_background_image()
create_initial()
