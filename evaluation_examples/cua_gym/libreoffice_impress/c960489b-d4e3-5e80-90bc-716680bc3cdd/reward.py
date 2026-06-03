"""
Reward Script: Change master slide to picture background with corporate_bg.jpg stretched to fill
Task ID: impress_ma_029
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Master slide background is PICTURE fill (not SOLID)
  Component 2 (0.3): Background image matches corporate_bg.jpg byte-for-byte
  Component 3 (0.3): Stretch fill mode used AND all 10 slides inherit from master
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_ma_029'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be a valid pptx (zip)
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        zf = zipfile.ZipFile(file_path, 'r')
    except Exception as e:
        print(f"CRITICAL: Cannot open pptx as zip: {e}")
        print("REWARD: 0.0")
        return 0.0

    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    ns_rels = 'http://schemas.openxmlformats.org/package/2006/relationships'

    # ---- Component 1: Master slide background is PICTURE fill (0.4 points) ----
    try:
        with zf.open('ppt/slideMasters/slideMaster1.xml') as f:
            master_root = ET.parse(f).getroot()

        bg = master_root.find(f'.//{{{ns_p}}}cSld/{{{ns_p}}}bg', )
        if bg is None:
            bg = master_root.find(f'.//{{{ns_p}}}bg')

        blip_fill = None
        if bg is not None:
            # Look for blipFill inside bgPr
            bgPr = bg.find(f'{{{ns_p}}}bgPr')
            if bgPr is not None:
                blip_fill = bgPr.find(f'{{{ns_a}}}blipFill')

        if blip_fill is not None:
            print(f"PASS: Component 1 — Master slide has picture (blipFill) background (0.4 pts)")
            total_score += 0.4
        else:
            # Check if it's still solid fill
            solid_fill = None
            if bg is not None and bgPr is not None:
                solid_fill = bgPr.find(f'{{{ns_a}}}solidFill')
            if solid_fill is not None:
                print(f"FAIL: Component 1 — Master slide still has solid fill background, expected picture")
            else:
                print(f"FAIL: Component 1 — Master slide background is not a picture fill")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---- Component 2: Background image matches corporate_bg.jpg (0.3 points) ----
    try:
        if blip_fill is not None:
            # Get the relationship ID from the blip element
            blip = blip_fill.find(f'{{{ns_a}}}blip')
            if blip is not None:
                embed_rid = blip.get(f'{{{ns_r}}}embed')
                if embed_rid:
                    # Resolve the relationship to find the image path
                    with zf.open('ppt/slideMasters/_rels/slideMaster1.xml.rels') as rf:
                        rels_root = ET.parse(rf).getroot()

                    image_target = None
                    for rel in rels_root:
                        if rel.get('Id') == embed_rid:
                            image_target = rel.get('Target')
                            break

                    if image_target:
                        # Target is relative (e.g., ../media/image1.jpg), resolve to zip path
                        if image_target.startswith('..'):
                            image_zip_path = 'ppt/' + image_target.lstrip('../')
                        else:
                            image_zip_path = 'ppt/slideMasters/' + image_target
                        # Normalize path
                        image_zip_path = os.path.normpath(image_zip_path).replace('\\', '/')

                        # Read the embedded image bytes
                        embedded_bytes = zf.read(image_zip_path)

                        # Read the source corporate_bg.jpg
                        source_path = os.path.join(WORKDIR, 'Desktop', 'corporate_bg.jpg')
                        if os.path.exists(source_path):
                            with open(source_path, 'rb') as sf:
                                source_bytes = sf.read()

                            if embedded_bytes == source_bytes:
                                print(f"PASS: Component 2 — Embedded image matches corporate_bg.jpg ({len(embedded_bytes)} bytes) (0.3 pts)")
                                total_score += 0.3
                            else:
                                print(f"FAIL: Component 2 — Embedded image ({len(embedded_bytes)} bytes) does not match corporate_bg.jpg ({len(source_bytes)} bytes)")
                        else:
                            print(f"FAIL: Component 2 — Source file corporate_bg.jpg not found at {source_path}")
                    else:
                        print(f"FAIL: Component 2 — Could not resolve image relationship {embed_rid}")
                else:
                    print(f"FAIL: Component 2 — No embed attribute on blip element")
            else:
                print(f"FAIL: Component 2 — No blip element found in blipFill")
        else:
            print(f"FAIL: Component 2 — No blipFill on master (depends on Component 1)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---- Component 3: Stretch fill mode AND all 10 slides inherit from master (0.3 points) ----
    try:
        sub_score = 0.0

        # Check stretch mode
        if blip_fill is not None:
            stretch = blip_fill.find(f'{{{ns_a}}}stretch')
            if stretch is not None:
                fill_rect = stretch.find(f'{{{ns_a}}}fillRect')
                if fill_rect is not None:
                    print(f"  Stretch fill with fillRect confirmed")
                    sub_score += 0.15
                else:
                    print(f"  Stretch element found but no fillRect")
                    sub_score += 0.05
            else:
                # Check for tile or other fill modes
                print(f"  No stretch element — image may not be stretched to fill")
        else:
            print(f"  No blipFill — cannot check stretch mode")

        # Check all slides still inherit from master (type 5 = BACKGROUND)
        # Only award points if blipFill is present (task-introduced change gate)
        if blip_fill is not None:
            from pptx import Presentation as PptxPresentation
            prs = PptxPresentation(file_path)
            slide_count = len(prs.slides)
            inherit_count = 0
            for i, slide in enumerate(prs.slides):
                fill_type = slide.background.fill.type
                if fill_type is not None and fill_type == 5:  # BACKGROUND = inherited
                    inherit_count += 1
                elif fill_type is None:
                    inherit_count += 1
                else:
                    print(f"  Slide {i} has override background (fill type {fill_type})")

            if slide_count >= 10 and inherit_count == slide_count:
                print(f"  All {slide_count} slides inherit master picture background")
                sub_score += 0.15
            elif inherit_count > 0:
                print(f"  {inherit_count}/{slide_count} slides inherit master background")
                sub_score += 0.15 * (inherit_count / slide_count)
            else:
                print(f"  No slides inherit master background")
        else:
            print(f"  Skipped inheritance check — no picture background on master")

        if sub_score > 0:
            print(f"PASS: Component 3 — Stretch + inheritance ({sub_score:.2f} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 3 — Neither stretch fill nor inheritance verified")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    zf.close()

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook: save any unsaved LibreOffice state
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
persist_app_state()
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
