"""
FINAL REWARD SCRIPT - SUCCESS
Task: On the slide I’m editing right now, how do I apply an underline to every bit of text—and also change all of that text, tables included, to the exact Dark Red 2 shade (#8B0000)?
Generated: 2025-09-10 14:26:46
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

# -------------------------------------------------------------
# Reward Script: Verify that every text run (including those
# inside tables) on the current slide is BOTH underlined AND
# coloured Dark Red 2 (#8B0000).
# -------------------------------------------------------------
# Scoring Logic (progressive):
#   • Up to 0.5 pts for underline compliance (ratio based)
#   • Up to 0.5 pts for colour compliance  (ratio based)
#   • Returns a float in [0.0, 1.0]
#   • Returns 0.0 if no text runs are found (cannot verify)
# -------------------------------------------------------------

FILE_PATH = "/home/user/on_the_slide_im_editing_right_now_how_do_i_apply_an_underline_to_every_bit_of_textand_also_change_al_golden.pptx"
TARGET_RGB = (0x8B, 0x00, 0x00)  # Dark Red 2


def analyse_slide(slide):
    """Inspect all text runs (including table cells) in a slide."""
    total_runs = underline_ok = colour_ok = 0

    for shape in slide.shapes:
        # Standard text frames (titles, text boxes, etc.)
        if getattr(shape, "text_frame", None) is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    total_runs += 1
                    font = run.font
                    if font is not None:
                        # Underline check
                        if font.underline:
                            underline_ok += 1
                        # Colour check – ensure explicit RGB and compare
                        if font.color is not None and font.color.rgb is not None:
                            if tuple(font.color.rgb) == TARGET_RGB:
                                colour_ok += 1
        # Table cells
        if getattr(shape, "has_table", False):
            tbl = shape.table
            for row in tbl.rows:
                for cell in row.cells:
                    if cell.text_frame is None:
                        continue
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            total_runs += 1
                            font = run.font
                            if font is not None:
                                if font.underline:
                                    underline_ok += 1
                                if font.color is not None and font.color.rgb is not None:
                                    if tuple(font.color.rgb) == TARGET_RGB:
                                        colour_ok += 1
    return total_runs, underline_ok, colour_ok


def verify_task(file_path: str) -> float:
    print(f"Verifying presentation: {file_path}")

    # --- 0. File existence (prerequisite, no points) ---
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # --- 1. Attempt to open the PPTX (prerequisite, no points) ---
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Unable to open PPTX: {exc}")
        print("REWARD: 0.0")
        return 0.0

    # For the purpose of this task, the slide "being edited" is interpreted
    # as the first slide. (If there are more, we still aggregate across all.)
    total_runs = underline_runs = colour_runs = 0

    for idx, slide in enumerate(prs.slides, start=1):
        t, u, c = analyse_slide(slide)
        print(f"Slide {idx}: total={t}, underline_ok={u}, colour_ok={c}")
        total_runs += t
        underline_runs += u
        colour_runs += c

    if total_runs == 0:
        print("✗ No text runs detected – nothing to verify.")
        print("REWARD: 0.0")
        return 0.0

    # --- 2. Progressive Scoring ---
    underline_ratio = underline_runs / total_runs
    colour_ratio = colour_runs / total_runs

    underline_score = 0.5 * underline_ratio
    colour_score = 0.5 * colour_ratio

    final_score = round(min(underline_score + colour_score, 1.0), 2)

    # --- 3. Reporting ---
    print(f"Underline compliance: {underline_runs}/{total_runs} = {underline_ratio:.2%} -> {underline_score:.2f} pts")
    print(f"Colour   compliance: {colour_runs}/{total_runs} = {colour_ratio:.2%} -> {colour_score:.2f} pts")
    print(f"Total score (capped at 1.0): {final_score}")
    print(f"REWARD: {final_score}")

    return final_score


# ------------------ main execution ---------------------------
if __name__ == "__main__":
    verify_task(FILE_PATH)

