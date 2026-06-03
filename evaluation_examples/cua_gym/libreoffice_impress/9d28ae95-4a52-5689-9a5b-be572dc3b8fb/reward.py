"""
Reward Script: Onboarding presentation for new employees
Task ID: impress_wf_032
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): 14 slides exist
  Component 2 (0.10): Slide 1 has 'Welcome to Acme Corp' title
  Component 3 (0.10): Slide 2 has agenda with hyperlinks (ppaction://hlinksldjump)
  Component 4 (0.10): Slide 3 has timeline shapes (>=4 auto_shapes)
  Component 5 (0.10): Slide 4 has 3-column layout (3 auto_shapes)
  Component 6 (0.10): Slide 5 org chart with connectors (auto_shapes + lines)
  Component 7 (0.10): Slide 6 has checkbox-style bullets (☐ char)
  Component 8 (0.10): Slides 8 and 10 have tables
  Component 9 (0.10): Backgrounds are #FAFAFA on all slides
  Component 10 (0.05): Headers use #388E3C color
"""

import os
import zipfile
import xml.etree.ElementTree as ET


WORKDIR = '/home/user'
TASK_ID = 'impress_wf_032'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)
    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: Presentation has exactly 14 slides (0.15 points)
    try:
        if num_slides == 14:
            print(f"PASS: Component 1 — 14 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 14 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Gate: need at least 14 slides for remaining checks
    if num_slides < 14:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Component 2: Slide 1 has 'Welcome to Acme Corp' title text (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide1_text += para.text.strip().lower() + " "
        if "welcome to acme corp" in slide1_text:
            print(f"PASS: Component 2 — 'Welcome to Acme Corp' found on slide 1 (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — 'Welcome to Acme Corp' not found. Text: {slide1_text[:100]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has agenda with slide hyperlinks (0.10 points)
    # Hyperlinks should use ppaction://hlinksldjump action
    try:
        hlink_count = 0
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide2.xml') as f:
                content = f.read().decode()
                import re
                hlink_count = len(re.findall(r'ppaction://hlinksldjump', content))
        if hlink_count >= 3:
            print(f"PASS: Component 3 — Slide 2 has {hlink_count} slide hyperlinks (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — expected >=3 slide hyperlinks on slide 2, found {hlink_count}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has timeline with horizontal shapes (>=4 auto_shapes) (0.10 points)
    try:
        slide3 = prs.slides[2]
        auto_shapes_3 = []
        for s in slide3.shapes:
            if str(s.shape_type) == 'AUTO_SHAPE (1)':
                auto_shapes_3.append(s)
        # Need at least 4 timeline event shapes (excluding the line/bar)
        # Timeline should have year labels like 2005, 2010, etc.
        timeline_shapes = 0
        for s in auto_shapes_3:
            if s.has_text_frame:
                txt = s.text_frame.text.strip()
                if txt:
                    timeline_shapes += 1
        if timeline_shapes >= 4:
            print(f"PASS: Component 4 — Slide 3 has {timeline_shapes} timeline shapes (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 4 — expected >=4 timeline shapes on slide 3, found {timeline_shapes}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has 3-column layout (3 auto_shapes with text) (0.10 points)
    try:
        slide4 = prs.slides[3]
        col_shapes = []
        for s in slide4.shapes:
            if str(s.shape_type) == 'AUTO_SHAPE (1)' and s.has_text_frame:
                txt = s.text_frame.text.strip()
                if txt:
                    col_shapes.append(s)
        if len(col_shapes) == 3:
            print(f"PASS: Component 5 — Slide 4 has 3 column shapes (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — expected 3 column shapes on slide 4, found {len(col_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 org chart with connected shapes (auto_shapes + connectors) (0.10 points)
    try:
        slide5 = prs.slides[4]
        auto_count = 0
        connector_count = 0
        for s in slide5.shapes:
            stype = str(s.shape_type)
            if stype == 'AUTO_SHAPE (1)':
                auto_count += 1
            elif 'LINE' in stype or 'FREEFORM' in stype:
                connector_count += 1
        # Org chart needs multiple boxes and connectors
        if auto_count >= 4 and connector_count >= 3:
            print(f"PASS: Component 6 — Slide 5 org chart: {auto_count} boxes, {connector_count} connectors (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 6 — expected org chart (>=4 boxes, >=3 connectors), found {auto_count} boxes, {connector_count} connectors")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 6 has checkbox-style bullets (☐ or similar checkbox char) (0.10 points)
    try:
        checkbox_count = 0
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide6.xml') as f:
                root = ET.parse(f).getroot()
                for para in root.findall('.//a:p', ns):
                    pPr = para.find('a:pPr', ns)
                    if pPr is not None:
                        buChar = pPr.find('a:buChar', ns)
                        if buChar is not None:
                            char = buChar.get('char', '')
                            # Common checkbox characters
                            if char in ('\u2610', '\u2611', '\u2612', '\u2713', '\u2714', '\u2717', '\u2718', '\u25A1', '\u25A0', '\u25CB', '\u25CF'):
                                checkbox_count += 1
        if checkbox_count >= 3:
            print(f"PASS: Component 7 — Slide 6 has {checkbox_count} checkbox bullets (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 7 — expected >=3 checkbox bullets on slide 6, found {checkbox_count}")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slides 8 and 10 have tables (0.10 points)
    try:
        tables_found = 0
        for slide_idx in [7, 9]:  # 0-indexed: slide 8 and 10
            slide = prs.slides[slide_idx]
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables_found += 1
                    break
        if tables_found == 2:
            print(f"PASS: Component 8 — Tables found on slides 8 and 10 (0.10 pts)")
            total_score += 0.10
        elif tables_found == 1:
            print(f"PARTIAL: Component 8 — Table found on 1 of 2 expected slides (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 8 — expected tables on slides 8 and 10, found {tables_found}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: All slides have #FAFAFA background (0.10 points)
    try:
        fafafa_count = 0
        for i, slide in enumerate(prs.slides):
            bg = slide.background.fill
            if bg.type == 1:  # SOLID
                try:
                    rgb = str(bg.fore_color.rgb)
                    if rgb.upper() == 'FAFAFA':
                        fafafa_count += 1
                except:
                    pass
        # At least 12 of 14 slides should have FAFAFA background
        if fafafa_count >= 12:
            print(f"PASS: Component 9 — {fafafa_count}/14 slides have #FAFAFA background (0.10 pts)")
            total_score += 0.10
        elif fafafa_count >= 8:
            print(f"PARTIAL: Component 9 — {fafafa_count}/14 slides have #FAFAFA background (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 9 — only {fafafa_count}/14 slides have #FAFAFA background")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    # Component 10: Headers use #388E3C green color (0.05 points)
    # Count slides that have at least one run with #388E3C
    try:
        slides_with_green = 0
        for i, slide in enumerate(prs.slides):
            found_green = False
            for shape in slide.shapes:
                if found_green:
                    break
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if found_green:
                            break
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    rgb = str(run.font.color.rgb).upper()
                                    if rgb == '388E3C':
                                        found_green = True
                                        break
                            except:
                                pass
            if found_green:
                slides_with_green += 1
        # Expect green headers on most slides (at least 10 of 14)
        if slides_with_green >= 10:
            print(f"PASS: Component 10 — {slides_with_green}/14 slides have green (#388E3C) headers (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 10 — expected >=10 slides with green headers, found {slides_with_green}")
    except Exception as e:
        print(f"ERROR: Component 10 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    # Also check Desktop as task says to save there
    alt_path = f'{WORKDIR}/Desktop/Onboarding.pptx'
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print(f"File not found: {file_path}")
        print(f"Also checked: {alt_path}")
        print("REWARD: 0.0")
        file_path = None

if file_path:
    verify_task(file_path)
