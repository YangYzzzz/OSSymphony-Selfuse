"""
Initial Setup: Create a blank presentation for the onboarding task.
Task ID: impress_wf_032
Domain: libreoffice_impress

The agent will be presented with a blank presentation and asked to create
a 14-slide onboarding presentation for Acme Corp.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_032'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    # Create a blank presentation with one empty slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
