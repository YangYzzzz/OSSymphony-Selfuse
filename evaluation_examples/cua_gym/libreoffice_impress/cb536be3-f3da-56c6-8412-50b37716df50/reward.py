"""
Reward Script: Environmental Report Presentation from Env_Data.xlsx
Task ID: impress_wf_058
Domain: libreoffice_impress
Scoring:
  Component 1: File exists and has 10 slides (0.15)
  Component 2: Slide 1 title is 'Environmental Impact Report 2023' with #2E7D32 color (0.10)
  Component 3: Background color #C8E6C9 on slides (0.10)
  Component 4: Charts on slides 3-6 (0.20)
  Component 5: Slide 7 has gauge rectangles (partial fill) (0.15)
  Component 6: Slide 8 has a table with initiative data (0.10)
  Component 7: Slide 9 table has green (#2E7D32) percentage values in Change column (0.10)
  Component 8: Slide 10 has 2024 sustainability goals content (0.10)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_058'


def get_slide_background_rgb(slide):
    """Get the background RGB color of a slide, handling inheritance."""
    fill = slide.background.fill
    if fill.type == 1:  # solid
        return fill.fore_color.rgb
    elif fill.type == 5:  # inherited
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return master_fill.fore_color.rgb
        except Exception:
            pass
    return None


def get_text_color(shape):
    """Get the font color of the first non-empty run in a shape."""
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                try:
                    if run.font.color.type is not None:
                        return run.font.color.rgb
                except Exception:
                    pass
    return None


def get_all_text(slide):
    """Get all text from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts


def has_chart(slide):
    """Check if a slide contains a chart."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
    return False


def has_table(slide):
    """Check if a slide contains a table. Returns the table shape or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape
    return None


def has_rectangle_pair(slide):
    """Check if slide has at least 2 auto shapes (rectangles for gauge).
    Returns True if there are 2 rectangles where one is wider (background)
    and one is narrower (fill) at the same position."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rects.append(shape)
    if len(rects) < 2:
        return False
    # Check that at least two rectangles share top position (gauge pattern)
    # and one is narrower than the other (partial fill)
    for i in range(len(rects)):
        for j in range(i + 1, len(rects)):
            r1, r2 = rects[i], rects[j]
            if r1.top == r2.top and r1.left == r2.left:
                if r1.width != r2.width:
                    return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: File has exactly 10 slides (0.15 points)
    try:
        num_slides = len(prs.slides)
        if num_slides == 10:
            print(f"PASS: Component 1 — 10 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    if len(prs.slides) < 10:
        print(f"CRITICAL: Not enough slides to verify remaining components")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Slide 1 title text and heading color #2E7D32 (0.10 points)
    try:
        slide1 = prs.slides[0]
        texts = get_all_text(slide1)
        full_text = " ".join(texts).lower()
        title_found = "environmental impact report 2023" in full_text
        # Check heading color
        heading_color_ok = any(
            str(run.font.color.rgb) == "2E7D32"
            for shape in slide1.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs if "environmental impact report" in para.text.lower()
            for run in para.runs
            if run.font.color.type is not None
        )

        if title_found and heading_color_ok:
            print(f"PASS: Component 2 — Title text and #2E7D32 color confirmed (0.10 pts)")
            total_score += 0.10
        elif title_found:
            print(f"PARTIAL: Component 2 — Title found but color not #2E7D32 (0.05 pts)")
            total_score += 0.05
        else:
            print(f"FAIL: Component 2 — Title not found. Texts: {texts[:3]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Background color #C8E6C9 on at least 8 of 10 slides (0.10 points)
    try:
        bg_count = 0
        for slide in prs.slides:
            bg = get_slide_background_rgb(slide)
            if bg is not None and str(bg) == "C8E6C9":
                bg_count += 1
        if bg_count >= 8:
            print(f"PASS: Component 3 — {bg_count}/10 slides have #C8E6C9 background (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Only {bg_count}/10 slides have #C8E6C9 background")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Charts on slides 3, 4, 5, 6 (0.20 points — 0.05 each)
    try:
        chart_slides = [2, 3, 4, 5]  # 0-indexed
        chart_score = 0.0
        for idx in chart_slides:
            if has_chart(prs.slides[idx]):
                chart_score += 0.05
                print(f"  PASS: Slide {idx+1} has chart")
            else:
                print(f"  FAIL: Slide {idx+1} missing chart")
        if chart_score > 0:
            print(f"PASS: Component 4 — Charts found on chart slides ({chart_score} pts)")
            total_score += chart_score
        else:
            print(f"FAIL: Component 4 — No charts found on slides 3-6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 7 has gauge (rectangle pair with partial fill) (0.15 points)
    try:
        slide7 = prs.slides[6]
        if has_rectangle_pair(slide7):
            print(f"PASS: Component 5 — Gauge rectangle pair found on slide 7 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — No gauge rectangle pair found on slide 7")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 8 has a table with green initiatives (0.10 points)
    try:
        slide8 = prs.slides[7]
        tbl_shape = has_table(slide8)
        if tbl_shape is not None:
            tbl = tbl_shape.table
            # Verify it's the initiatives table — check header row
            headers = [tbl.cell(0, c).text.strip().lower() for c in range(len(tbl.columns))]
            has_initiative_col = any("initiative" in h or "name" in h for h in headers)
            has_status_col = any("status" in h for h in headers)
            if has_initiative_col and has_status_col and len(tbl.rows) >= 3:
                print(f"PASS: Component 6 — Initiatives table found ({len(tbl.rows)} rows, headers: {headers}) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 — Table found but not initiatives table. Headers: {headers}, rows: {len(tbl.rows)}")
        else:
            print(f"FAIL: Component 6 — No table found on slide 8")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 9 table has green (#2E7D32) percentage values in Change column (0.10 points)
    try:
        slide9 = prs.slides[8]
        tbl_shape = has_table(slide9)
        if tbl_shape is not None:
            tbl = tbl_shape.table
            # Find Change column index
            change_col = None
            for c in range(len(tbl.columns)):
                if "change" in tbl.cell(0, c).text.strip().lower():
                    change_col = c
                    break
            if change_col is not None:
                green_pct_count = 0
                for r in range(1, len(tbl.rows)):
                    cell = tbl.cell(r, change_col)
                    cell_text = cell.text.strip()
                    if "%" in cell_text:
                        # Check if color is green
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                try:
                                    if run.font.color.type is not None:
                                        if str(run.font.color.rgb) == "2E7D32":
                                            green_pct_count += 1
                                except Exception:
                                    pass
                if green_pct_count >= 3:
                    print(f"PASS: Component 7 — {green_pct_count} green percentage values found in Change column (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 7 — Only {green_pct_count} green percentage values in Change column (need >= 3)")
            else:
                print(f"FAIL: Component 7 — No 'Change' column found in slide 9 table")
        else:
            print(f"FAIL: Component 7 — No table found on slide 9")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 10 has 2024 sustainability goals content (0.10 points)
    try:
        slide10 = prs.slides[9]
        texts = get_all_text(slide10)
        full_text = " ".join(texts).lower()
        has_goals_title = "2024" in full_text and ("goal" in full_text or "sustainability" in full_text)
        # Must have multiple goal items (at least 3 paragraphs with content)
        goal_items = 0
        for shape in slide10.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if len(t) > 20:  # goal descriptions are longer than 20 chars
                        goal_items += 1
        if has_goals_title and goal_items >= 3:
            print(f"PASS: Component 8 — 2024 goals slide found with {goal_items} goal items (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 8 — has_goals_title={has_goals_title}, goal_items={goal_items}")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/Desktop/Environmental_Report.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
