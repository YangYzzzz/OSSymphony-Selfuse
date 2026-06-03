"""
Reward Script: Create a data-rich dashboard slide (slide 4) with four mini-charts in 2x2 grid
Task ID: impress_gf2_034
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Bar chart present on slide 4 with COLUMN type
  Component 2 (0.15): Bar chart data matches (Jan=45, Feb=52, Mar=48)
  Component 3 (0.15): Pie chart present on slide 4 with PIE type
  Component 4 (0.15): Pie chart data matches (A=40%, B=35%, C=25%)
  Component 5 (0.10): Line chart present on slide 4 with LINE type
  Component 6 (0.10): Line chart data matches (42, 45, 51, 48)
  Component 7 (0.10): Table present on slide 4 with 3x2 dimensions
  Component 8 (0.10): Table has KPI names and values
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_034'


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice changes."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing python-pptx library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect charts and tables on slide 4
    charts = []
    tables = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            charts.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tables.append(shape)

    print(f"INFO: Slide 4 has {len(charts)} chart(s) and {len(tables)} table(s)")

    # Classify charts by type
    bar_charts = []
    pie_charts = []
    line_charts = []

    for shape in charts:
        chart = shape.chart
        ct = chart.chart_type
        # COLUMN_CLUSTERED=51, BAR_CLUSTERED=57, etc.
        if ct in (51, 52, 53, 54, 55, 56, 57, 58, 59, 60):
            bar_charts.append(shape)
        # PIE=5, PIE_EXPLODED=69, PIE_OF_PIE=68, etc.
        elif ct in (5, 68, 69, 70, 71, 72):
            pie_charts.append(shape)
        # LINE=4, LINE_MARKERS=65, LINE_STACKED=63, etc.
        elif ct in (4, 63, 64, 65, 66, 67):
            line_charts.append(shape)

    # Component 1: Bar chart exists on slide 4 (0.15 points)
    try:
        if len(bar_charts) >= 1:
            print(f"PASS: Component 1 — Bar/column chart found on slide 4 (type={bar_charts[0].chart.chart_type}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — No bar/column chart found on slide 4. Chart types found: {[s.chart.chart_type for s in charts]}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Bar chart has correct data (Jan=45, Feb=52, Mar=48) (0.15 points)
    try:
        if len(bar_charts) >= 1:
            chart = bar_charts[0].chart
            plot = chart.plots[0]
            values = list(plot.series[0].values)
            categories = list(plot.categories)

            expected_values = [45.0, 52.0, 48.0]
            expected_cats_keywords = ['jan', 'feb', 'mar']

            cats_lower = [str(c).lower() for c in categories]
            cats_match = len(cats_lower) == 3 and all(
                any(kw in cat for cat in cats_lower) for kw in expected_cats_keywords
            )
            vals_match = len(values) == 3 and all(
                abs(values[i] - expected_values[i]) < 0.01 for i in range(3)
            )

            if cats_match and vals_match:
                print(f"PASS: Component 2 — Bar chart data correct: cats={categories}, vals={values} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 2 — Bar chart data mismatch: cats={categories} (expect Jan/Feb/Mar), vals={values} (expect {expected_values})")
        else:
            print(f"FAIL: Component 2 — No bar chart to check data")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Pie chart exists on slide 4 (0.15 points)
    try:
        if len(pie_charts) >= 1:
            print(f"PASS: Component 3 — Pie chart found on slide 4 (type={pie_charts[0].chart.chart_type}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — No pie chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Pie chart has correct data (A=40%, B=35%, C=25%) (0.15 points)
    try:
        if len(pie_charts) >= 1:
            chart = pie_charts[0].chart
            plot = chart.plots[0]
            values = list(plot.series[0].values)
            categories = list(plot.categories)

            expected_values = [40.0, 35.0, 25.0]
            # Categories should mention Product A/B/C or just A/B/C
            cats_lower = [str(c).lower() for c in categories]

            vals_match = len(values) == 3 and all(
                abs(values[i] - expected_values[i]) < 0.01 for i in range(3)
            )
            cats_match = len(cats_lower) == 3 and all(
                any(letter in cat for cat in cats_lower)
                for letter in ['a', 'b', 'c']
            )

            if cats_match and vals_match:
                print(f"PASS: Component 4 — Pie chart data correct: cats={categories}, vals={values} (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 — Pie chart data mismatch: cats={categories}, vals={values} (expect {expected_values})")
        else:
            print(f"FAIL: Component 4 — No pie chart to check data")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Line chart exists on slide 4 (0.10 points)
    try:
        if len(line_charts) >= 1:
            print(f"PASS: Component 5 — Line chart found on slide 4 (type={line_charts[0].chart.chart_type}) (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 5 — No line chart found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Line chart has correct NPS data (42, 45, 51, 48) (0.10 points)
    try:
        if len(line_charts) >= 1:
            chart = line_charts[0].chart
            plot = chart.plots[0]
            values = list(plot.series[0].values)

            expected_values = [42.0, 45.0, 51.0, 48.0]
            vals_match = len(values) == 4 and all(
                abs(values[i] - expected_values[i]) < 0.01 for i in range(4)
            )

            if vals_match:
                print(f"PASS: Component 6 — Line chart data correct: vals={values} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 6 — Line chart data mismatch: vals={values} (expect {expected_values})")
        else:
            print(f"FAIL: Component 6 — No line chart to check data")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Table exists on slide 4 with approximately 3 rows and 2 columns (0.10 points)
    try:
        if len(tables) >= 1:
            table = tables[0].table
            nrows = len(table.rows)
            ncols = len(table.columns)
            # Task says 2x3 metric table with KPI names and values -> 3 rows x 2 cols (or 2 cols x 3 rows)
            if nrows >= 2 and ncols >= 2:
                print(f"PASS: Component 7 — Table found on slide 4: {nrows}x{ncols} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 7 — Table dimensions too small: {nrows}x{ncols}")
        else:
            print(f"FAIL: Component 7 — No table found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Table contains KPI-style content (names + values) (0.10 points)
    try:
        if len(tables) >= 1:
            table = tables[0].table
            nrows = len(table.rows)
            ncols = len(table.columns)
            # Collect all cell text
            all_text = []
            for r in range(nrows):
                for c in range(ncols):
                    all_text.append(table.cell(r, c).text.strip().lower())
            all_joined = ' '.join(all_text)

            # Check that table has at least some KPI-style content
            # We expect metric names and numeric-ish values
            has_names = False
            has_values = False

            # Check for KPI-like keywords (revenue, growth, nps, etc.)
            kpi_keywords = ['revenue', 'growth', 'nps', 'kpi', 'metric', 'customer', 'rate', 'score', 'sales']
            for kw in kpi_keywords:
                if kw in all_joined:
                    has_names = True
                    break

            # Check for numeric-like values
            import re
            for txt in all_text:
                if re.search(r'\d', txt):
                    has_values = True
                    break

            if has_names and has_values:
                print(f"PASS: Component 8 — Table has KPI names and values (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 8 — Table missing KPI content. has_names={has_names}, has_values={has_values}. Content: {all_text}")
        else:
            print(f"FAIL: Component 8 — No table to check content")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
