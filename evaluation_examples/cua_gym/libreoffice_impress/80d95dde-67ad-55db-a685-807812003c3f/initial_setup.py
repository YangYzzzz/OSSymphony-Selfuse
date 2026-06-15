"""
Initial Setup: Create a 6-slide Management Dashboard presentation with an empty Slide 4
Task ID: impress_gf2_034
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_034'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and add body text to a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find the body placeholder (index 1 typically)
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break
    if body_ph and body_lines:
        tf = body_ph.text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q1 2025 Management Dashboard"
    slide1.placeholders[1].text = "Quarterly Business Review\nPrepared by Strategy & Analytics Team"

    # --- Slide 2: Revenue Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Revenue Overview", [
        "Total Q1 Revenue: $4.2M (+12% YoY)",
        "North America: $2.1M (50%)",
        "EMEA: $1.3M (31%)",
        "APAC: $0.8M (19%)",
        "Gross margin improved to 68.5% from 65.2%",
        "Enterprise segment grew 23% driven by new logos",
    ])

    # --- Slide 3: Team Performance ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Team Performance", [
        "Engineering: 94% sprint completion rate",
        "Sales: 108% of Q1 quota achieved",
        "Customer Success: NPS improved from 42 to 51",
        "Marketing: 2,340 MQLs generated (+18% QoQ)",
        "New hires onboarded: 14 across 3 departments",
        "Employee satisfaction score: 4.3/5.0",
    ])

    # --- Slide 4: Performance Dashboard (EMPTY - task target) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Performance Dashboard"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)
    # Content area is completely empty - no charts, no tables

    # --- Slide 5: Customer Feedback ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Customer Feedback", [
        "Overall satisfaction: 4.5/5.0 (survey of 1,200 customers)",
        "Top praise: Onboarding experience and support response times",
        "Key concern: Mobile app performance on Android devices",
        "Feature requests: Advanced reporting, API integrations",
        "Churn rate decreased to 2.1% from 3.4% in Q4",
        "Reference customers increased to 45 (+8 this quarter)",
    ])

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Next Steps & Q2 Priorities", [
        "Launch mobile app v2.0 by April 30",
        "Expand APAC sales team with 3 new AEs",
        "Complete SOC 2 Type II certification",
        "Roll out advanced analytics dashboard for enterprise tier",
        "Host annual customer summit in June",
        "Target: $4.8M Q2 revenue (15% growth)",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
