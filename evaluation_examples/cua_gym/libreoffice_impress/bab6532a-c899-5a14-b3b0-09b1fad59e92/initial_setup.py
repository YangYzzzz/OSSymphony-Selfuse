"""
Initial Setup: Strategic Plan presentation with blank SWOT Analysis slide
Task ID: impress_exec_007
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_007'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Strategic Plan 2025-2027"
    slide1.placeholders[1].text = "Meridian Technologies Inc.\nBoard Presentation — Q2 2025"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Company Overview & Mission"
    items = [
        "Market Analysis & Competitive Landscape",
        "Financial Performance Review",
        "SWOT Analysis",
        "Strategic Initiatives for 2025-2027",
        "Resource Allocation & Budget",
        "Risk Management Framework",
        "Implementation Timeline",
        "Q&A",
    ]
    for item in items:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Company Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Company Overview"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Founded in 2012, Meridian Technologies is a mid-market SaaS provider"
    overview_items = [
        "Revenue: $127.4M (FY2024), up 23% YoY",
        "Employees: 845 across 6 offices globally",
        "Products: CloudSync Pro, DataBridge, InsightHub",
        "Markets: North America (62%), Europe (28%), APAC (10%)",
        "Net Promoter Score: 72 (industry avg: 41)",
    ]
    for item in overview_items:
        p = body3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Financial Performance ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Financial Performance FY2024"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Key Financial Highlights"
    finance_items = [
        "Total Revenue: $127.4M (+23% YoY)",
        "Gross Margin: 71.3% (up from 68.9%)",
        "Operating Income: $18.2M (+31% YoY)",
        "ARR: $142.8M with 115% net dollar retention",
        "Free Cash Flow: $22.7M (17.8% FCF margin)",
        "Customer Count: 2,340 enterprise accounts",
    ]
    for item in finance_items:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: SWOT Analysis (BLANK - task target) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title textbox at the top
    txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "SWOT Analysis"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 6: Strategic Initiatives ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Strategic Initiatives 2025-2027"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Initiative 1: AI-Powered Analytics Platform"
    strat_items = [
        "Initiative 2: European Market Expansion (UK, DACH, Nordics)",
        "Initiative 3: Enterprise Tier Launch ($250K+ ACV)",
        "Initiative 4: Strategic Partnerships with AWS & Salesforce",
        "Initiative 5: Developer Platform & API Marketplace",
    ]
    for item in strat_items:
        p = body6.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 7: Resource Allocation ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Resource Allocation & Budget"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "R&D Investment: $38.2M (30% of revenue)"
    budget_items = [
        "Sales & Marketing: $35.8M (28% of revenue)",
        "G&A: $12.7M (10% of revenue)",
        "New Headcount: 120 positions planned for FY2025",
        "Capital Expenditure: $8.5M (infrastructure upgrades)",
        "M&A Reserve: $15M for strategic acquisitions",
    ]
    for item in budget_items:
        p = body7.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 8: Risk Management ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Risk Management Framework"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Market Risk: Increased competition from enterprise players"
    risk_items = [
        "Technology Risk: AI model reliability and data privacy compliance",
        "Operational Risk: Scaling team across geographies",
        "Financial Risk: Currency exposure from EU expansion",
        "Regulatory Risk: GDPR, SOC2, and emerging AI regulations",
        "Mitigation: Quarterly risk reviews with executive committee",
    ]
    for item in risk_items:
        p = body8.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 9: Timeline ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Implementation Timeline"
    body9 = slide9.placeholders[1].text_frame
    body9.text = "Q2 2025: AI analytics beta launch & UK office setup"
    timeline_items = [
        "Q3 2025: Enterprise tier pilot with 10 accounts",
        "Q4 2025: AWS partnership go-live, DACH market entry",
        "Q1 2026: Full AI platform GA, developer API beta",
        "Q2 2026: Nordics expansion, API marketplace launch",
        "Q4 2026: Series D fundraise ($80M target)",
    ]
    for item in timeline_items:
        p = body9.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 10: Q&A ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox10 = slide10.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
    tf10 = txBox10.text_frame
    p10 = tf10.paragraphs[0]
    p10.text = "Questions & Discussion"
    p10.alignment = PP_ALIGN.CENTER
    run10 = p10.runs[0]
    run10.font.size = Pt(44)
    run10.font.bold = True
    run10.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    p10b = tf10.add_paragraph()
    p10b.text = "Thank you for your time"
    p10b.alignment = PP_ALIGN.CENTER
    run10b = p10b.runs[0]
    run10b.font.size = Pt(20)
    run10b.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
