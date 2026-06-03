"""
Reward Script: SWOT Analysis 2x2 grid on slide 5
Task ID: impress_exec_007
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Four rectangles with correct SWOT labels exist on slide 5
  Component 2 (0.35): Correct fill colors (#4CAF50, #F44336, #2196F3, #FF9800)
  Component 3 (0.25): Correct rectangle sizes (~6in wide x ~3in tall)
  Component 4 (0.15): Correct 2x2 grid layout positions
"""

import os

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_007'

# Expected SWOT entries: label -> fill color hex string
EXPECTED_RECTS = {
    'Strengths':     '4CAF50',
    'Weaknesses':    'F44336',
    'Opportunities': '2196F3',
    'Threats':       'FF9800',
}

# Approximate target size in EMU (6 inches wide, 3 inches tall)
TARGET_WIDTH = Inches(6)   # 5486400 EMU
TARGET_HEIGHT = Inches(3)  # 2743200 EMU
SIZE_TOLERANCE = 0.15  # 15% relative tolerance for size


def is_approx(actual, expected, tol=SIZE_TOLERANCE):
    """Check if actual is within tol relative tolerance of expected."""
    if expected == 0:
        return actual == 0
    return abs(actual - expected) / expected <= tol


def find_swot_rectangles(slide):
    """
    Find rectangle-like shapes on the slide that contain SWOT labels.
    Returns dict: label -> shape, for labels found in EXPECTED_RECTS.
    """
    found = {}
    for shape in slide.shapes:
        # Accept AUTO_SHAPE (rectangles, rounded rects, etc.)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            text = shape.text.strip() if hasattr(shape, 'text') else ''
            for label in EXPECTED_RECTS:
                if label.lower() == text.lower():
                    found[label] = shape
                    break
    return found


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # 0-indexed, slide 5
    rects = find_swot_rectangles(slide)

    # Component 1: Four rectangles with correct SWOT labels (0.25 points)
    try:
        found_labels = set(rects.keys())
        expected_labels = set(EXPECTED_RECTS.keys())
        matching = found_labels & expected_labels
        count = len(matching)

        if count == 4:
            print(f"PASS: Component 1 — All 4 SWOT rectangles found: {sorted(matching)} (0.25 pts)")
            total_score += 0.25
        elif count > 0:
            partial = 0.25 * (count / 4)
            print(f"PARTIAL: Component 1 — {count}/4 rectangles found: {sorted(matching)} ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 — No SWOT rectangles found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Correct fill colors (0.35 points, ~0.0875 per rect)
    try:
        color_matches = 0
        for label, expected_color in EXPECTED_RECTS.items():
            if label not in rects:
                print(f"FAIL: Component 2 — '{label}' rectangle not found, cannot check color")
                continue
            shape = rects[label]
            try:
                fill = shape.fill
                if fill.type == 1:  # SOLID fill
                    actual_color = str(fill.fore_color.rgb).upper()
                    expected_upper = expected_color.upper()
                    if actual_color == expected_upper:
                        print(f"PASS: Component 2 — '{label}' fill color correct: #{actual_color}")
                        color_matches += 1
                    else:
                        print(f"FAIL: Component 2 — '{label}' fill color expected #{expected_upper}, found #{actual_color}")
                else:
                    print(f"FAIL: Component 2 — '{label}' fill is not solid (type={fill.type})")
            except Exception as e:
                print(f"ERROR: Component 2 — '{label}' fill check failed: {e}")

        if color_matches > 0:
            color_score = 0.35 * (color_matches / 4)
            print(f"Component 2 subtotal: {color_matches}/4 colors correct ({color_score:.3f} pts)")
            total_score += color_score
        else:
            print(f"FAIL: Component 2 — No correct fill colors found")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct sizes ~6in x ~3in (0.25 points, ~0.0625 per rect)
    try:
        size_matches = 0
        for label in EXPECTED_RECTS:
            if label not in rects:
                print(f"FAIL: Component 3 — '{label}' rectangle not found, cannot check size")
                continue
            shape = rects[label]
            w_ok = is_approx(shape.width, TARGET_WIDTH)
            h_ok = is_approx(shape.height, TARGET_HEIGHT)
            if w_ok and h_ok:
                print(f"PASS: Component 3 — '{label}' size correct: {shape.width/914400:.2f}in x {shape.height/914400:.2f}in")
                size_matches += 1
            else:
                print(f"FAIL: Component 3 — '{label}' size: {shape.width/914400:.2f}in x {shape.height/914400:.2f}in (expected ~6.00in x ~3.00in)")

        if size_matches > 0:
            size_score = 0.25 * (size_matches / 4)
            print(f"Component 3 subtotal: {size_matches}/4 sizes correct ({size_score:.3f} pts)")
            total_score += size_score
        else:
            print(f"FAIL: Component 3 — No rectangles with correct sizes")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 2x2 grid layout (0.15 points)
    # Verify: top-left / top-right share a top position; bottom-left / bottom-right share a top;
    #         left pair share left position; right pair share left position;
    #         top pair above bottom pair; left pair left of right pair.
    try:
        required = ['Strengths', 'Weaknesses', 'Opportunities', 'Threats']
        if all(r in rects for r in required):
            s = rects['Strengths']
            w = rects['Weaknesses']
            o = rects['Opportunities']
            t = rects['Threats']

            layout_checks = 0

            # Check: Strengths is left of Weaknesses (top row)
            if s.left < w.left:
                layout_checks += 1
            else:
                print(f"FAIL: Component 4 — Strengths ({s.left}) should be left of Weaknesses ({w.left})")

            # Check: Opportunities is left of Threats (bottom row)
            if o.left < t.left:
                layout_checks += 1
            else:
                print(f"FAIL: Component 4 — Opportunities ({o.left}) should be left of Threats ({t.left})")

            # Check: Strengths is above Opportunities (left column)
            if s.top < o.top:
                layout_checks += 1
            else:
                print(f"FAIL: Component 4 — Strengths ({s.top}) should be above Opportunities ({o.top})")

            # Check: Weaknesses is above Threats (right column)
            if w.top < t.top:
                layout_checks += 1
            else:
                print(f"FAIL: Component 4 — Weaknesses ({w.top}) should be above Threats ({t.top})")

            if layout_checks == 4:
                print(f"PASS: Component 4 — Correct 2x2 grid layout (0.15 pts)")
                total_score += 0.15
            elif layout_checks > 0:
                partial = 0.15 * (layout_checks / 4)
                print(f"PARTIAL: Component 4 — {layout_checks}/4 layout checks passed ({partial:.3f} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 4 — Grid layout incorrect")
        else:
            missing = [r for r in required if r not in rects]
            print(f"FAIL: Component 4 — Missing rectangles for layout check: {missing}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
