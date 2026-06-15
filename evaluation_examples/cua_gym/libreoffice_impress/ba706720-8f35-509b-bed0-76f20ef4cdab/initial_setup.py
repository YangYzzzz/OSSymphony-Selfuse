"""
Initial Setup: Create a 6-slide Lab Report presentation with slide 4 having
a title and single content area.
Task ID: impress_teach_081
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_081'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, alignment=PP_ALIGN.LEFT):
    """Helper to add a simple text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return txBox


def add_title_to_slide(slide, title_text):
    """Add a title text box at the top of a blank slide."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x3C, 0x6E)
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Title Only or Blank

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Organic Chemistry Lab Report"
    slide1.placeholders[1].text = "Dr. Elena Vasquez | CHEM 342 | Spring 2025"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(blank_layout)
    add_title_to_slide(slide2, "Introduction")
    add_text_box(slide2, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 "This experiment investigates the Fischer esterification reaction "
                 "between acetic acid and 1-butanol in the presence of sulfuric acid "
                 "as a catalyst. The objective is to synthesize n-butyl acetate and "
                 "measure the reaction yield under varying temperature conditions.\n\n"
                 "Fischer esterification is a classic organic reaction that produces "
                 "an ester from a carboxylic acid and an alcohol. Understanding "
                 "the kinetics and equilibrium of this reaction is essential for "
                 "industrial chemical synthesis and pharmaceutical applications.",
                 font_size=14)

    # --- Slide 3: Materials and Methods ---
    slide3 = prs.slides.add_slide(blank_layout)
    add_title_to_slide(slide3, "Materials and Methods")
    add_text_box(slide3, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 "Reagents:\n"
                 "  - Glacial acetic acid (17.4 M, 15 mL)\n"
                 "  - 1-Butanol (10.9 M, 20 mL)\n"
                 "  - Concentrated sulfuric acid (1 mL, catalyst)\n"
                 "  - Sodium bicarbonate solution (5% w/v, 50 mL)\n\n"
                 "Equipment:\n"
                 "  - 250 mL round-bottom flask with reflux condenser\n"
                 "  - Heating mantle with temperature controller\n"
                 "  - Separatory funnel (125 mL)\n"
                 "  - Rotary evaporator\n"
                 "  - IR spectrometer (Bruker Alpha II)",
                 font_size=13)

    # --- Slide 4: Experiment Setup (single content area — task target) ---
    slide4 = prs.slides.add_slide(blank_layout)
    add_title_to_slide(slide4, "Experiment Setup")
    add_text_box(slide4, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 "The reflux apparatus was assembled with a 250 mL round-bottom flask "
                 "connected to a water-cooled condenser. Acetic acid and 1-butanol "
                 "were mixed in the flask, and sulfuric acid was added dropwise as "
                 "the catalyst. The reaction mixture was heated to reflux temperature "
                 "(approximately 118°C) and maintained for 60 minutes with continuous "
                 "stirring using a magnetic stir bar.",
                 font_size=14)

    # --- Slide 5: Results ---
    slide5 = prs.slides.add_slide(blank_layout)
    add_title_to_slide(slide5, "Results")
    add_text_box(slide5, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 "Trial 1 (100°C): Yield 42.3%, reaction time 90 min\n"
                 "Trial 2 (110°C): Yield 58.7%, reaction time 75 min\n"
                 "Trial 3 (118°C): Yield 71.2%, reaction time 60 min\n"
                 "Trial 4 (125°C): Yield 68.9%, reaction time 45 min\n\n"
                 "The IR spectrum of the purified product showed a strong carbonyl "
                 "stretch at 1742 cm⁻¹, confirming ester formation. No broad O-H "
                 "stretch was observed, indicating complete removal of unreacted "
                 "starting materials during purification.",
                 font_size=14)

    # --- Slide 6: Conclusion ---
    slide6 = prs.slides.add_slide(blank_layout)
    add_title_to_slide(slide6, "Conclusion")
    add_text_box(slide6, Inches(0.5), Inches(1.3), Inches(9), Inches(5.5),
                 "The Fischer esterification of acetic acid with 1-butanol "
                 "achieved an optimal yield of 71.2% at 118°C. Higher temperatures "
                 "did not improve yield due to competing side reactions. The "
                 "sulfuric acid catalyst concentration of 5% v/v proved effective "
                 "without causing significant charring.\n\n"
                 "Future work should explore alternative catalysts such as "
                 "p-toluenesulfonic acid and investigate the effect of reactant "
                 "molar ratios on equilibrium conversion.",
                 font_size=14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
