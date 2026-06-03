"""
Reward Script: Two-column layout on slide 4 with image placeholder and bulleted list
Task ID: impress_teach_081
Domain: libreoffice_impress
Scoring:
  Component 1: Rectangle placeholder exists with ~5x4 inch size (0.30 pts)
  Component 2: Rectangle has dashed border, #F5F5F5 fill, and 'Insert Diagram' text (0.35 pts)
  Component 3: Right-side text box with bulleted list items (0.35 pts)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_081'


def is_approx(val1, val2, tolerance=0.1):
    """Check if two values are approximately equal (within tolerance ratio)."""
    if val1 == val2:
        return True
    if val1 == 0 or val2 == 0:
        return abs(val1 - val2) < 0.05
    return abs(val1 - val2) / max(abs(val1), abs(val2)) <= tolerance


def persist_app_state(domain):
    """Save any unsaved LibreOffice state."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)
    shapes = list(slide.shapes)
    EMU_PER_INCH = 914400

    # Find candidate rectangle shapes (AUTO_SHAPE type) on slide 4
    rect_shapes = []
    textbox_shapes = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rect_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            textbox_shapes.append(shape)

    # =========================================================
    # Component 1: Rectangle placeholder exists with ~5x4 inch size (0.30 pts)
    # This checks for a new AUTO_SHAPE rectangle with approximate 5x4 inch dimensions.
    # Initial state has 0 AUTO_SHAPEs on slide 4, so this only passes on golden.
    # =========================================================
    target_rect = None
    try:
        for shape in rect_shapes:
            w_inches = shape.width / EMU_PER_INCH
            h_inches = shape.height / EMU_PER_INCH
            if is_approx(w_inches, 5.0, 0.15) and is_approx(h_inches, 4.0, 0.15):
                target_rect = shape
                break

        if target_rect is not None:
            w_in = target_rect.width / EMU_PER_INCH
            h_in = target_rect.height / EMU_PER_INCH
            print(f"PASS: Component 1 -- Rectangle found, size {w_in:.2f}x{h_in:.2f} in (0.30 pts)")
            total_score += 0.30
        else:
            if rect_shapes:
                for s in rect_shapes:
                    print(f"  Found rect: {s.width/EMU_PER_INCH:.2f}x{s.height/EMU_PER_INCH:.2f} in")
            print(f"FAIL: Component 1 -- No rectangle ~5x4 inches found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================
    # Component 2: Rectangle has dashed border, #F5F5F5 fill, 'Insert Diagram' text (0.35 pts)
    # Sub-checks: fill color (0.1), dashed border (0.1), text content (0.15)
    # =========================================================
    try:
        if target_rect is not None:
            comp2_score = 0.0

            # Sub-check 2a: Fill color #F5F5F5 (0.10 pts)
            try:
                fill = target_rect.fill
                if fill.type is not None and fill.type == 1:  # SOLID fill
                    rgb = str(fill.fore_color.rgb).upper()
                    if rgb == 'F5F5F5':
                        print(f"PASS: Component 2a -- Fill color is #F5F5F5 (0.10 pts)")
                        comp2_score += 0.10
                    else:
                        print(f"FAIL: Component 2a -- Fill color is #{rgb}, expected #F5F5F5")
                else:
                    print(f"FAIL: Component 2a -- Fill type is {fill.type}, expected SOLID (1)")
            except Exception as e:
                print(f"ERROR: Component 2a -- {e}")

            # Sub-check 2b: Dashed border (0.10 pts)
            try:
                line = target_rect.line
                dash = line.dash_style
                # Acceptable dashed styles: DASH (4), DASH_DOT (5), LONG_DASH (7), etc.
                # We accept any non-None, non-SOLID dash style
                if dash is not None and dash != 1:  # 1 = SOLID
                    print(f"PASS: Component 2b -- Dashed border found (dash_style={dash}) (0.10 pts)")
                    comp2_score += 0.10
                else:
                    print(f"FAIL: Component 2b -- Border dash_style={dash}, expected dashed")
            except Exception as e:
                print(f"ERROR: Component 2b -- {e}")

            # Sub-check 2c: Text 'Insert Diagram' (0.15 pts)
            try:
                shape_text = target_rect.text.strip()
                if 'Insert Diagram' in shape_text:
                    print(f"PASS: Component 2c -- Text 'Insert Diagram' found (0.15 pts)")
                    comp2_score += 0.15
                else:
                    print(f"FAIL: Component 2c -- Text is '{shape_text}', expected 'Insert Diagram'")
            except Exception as e:
                print(f"ERROR: Component 2c -- {e}")

            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 -- No target rectangle found, skipping sub-checks")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================
    # Component 3: Right-side text box with bulleted list (0.35 pts)
    # Must be a NEW text box (not the pre-existing title or content area)
    # positioned on the right side of the slide, containing bulleted items.
    # Initial state has only pre-existing text boxes, no bulleted list box.
    # Sub-checks: text box exists on right side (0.15), has bullets (0.20)
    # =========================================================
    try:
        # Look for a text box on the right half of the slide with multiple paragraphs
        # Pre-existing shapes: TextBox 1 (title) at left=457200 and TextBox 2 (content) at left=457200
        # Both span full width (8229600 EMU = 9in). New text box should be on the right side.
        slide_midpoint = prs.slide_width / 2  # ~5 inches = 4572000 EMU

        candidate_tb = None
        for shape in textbox_shapes:
            # Must be positioned on the right half
            if shape.left < slide_midpoint:
                continue
            # Must have at least 2 non-empty paragraphs (bulleted list)
            non_empty = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(non_empty) >= 2:
                candidate_tb = shape
                break

        comp3_score = 0.0

        if candidate_tb is not None:
            non_empty = [p for p in candidate_tb.text_frame.paragraphs if p.text.strip()]
            print(f"PASS: Component 3a -- Right-side text box found with {len(non_empty)} items (0.15 pts)")
            comp3_score += 0.15

            # Sub-check 3b: Paragraphs have bullet formatting
            # Check XML for bullet characters
            has_bullets = False
            try:
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                with zipfile.ZipFile(file_path, 'r') as zf:
                    with zf.open('ppt/slides/slide4.xml') as f:
                        root = ET.parse(f).getroot()
                        # Find TextBox 4 or any textbox shape matching our candidate
                        bullet_count = 0
                        total_paras = 0
                        for sp in root.findall('.//p:cSld//p:sp', ns):
                            cNvPr = sp.find('.//p:nvSpPr//p:cNvPr', ns)
                            sp_name = cNvPr.get('name') if cNvPr is not None else ''
                            if sp_name != candidate_tb.name:
                                continue
                            txBody = sp.find('.//p:txBody', ns)
                            if txBody is None:
                                continue
                            for para in txBody.findall('a:p', ns):
                                text = ''.join(t.text or '' for t in para.findall('.//a:t', ns))
                                if not text.strip():
                                    continue
                                total_paras += 1
                                pPr = para.find('a:pPr', ns)
                                if pPr is not None:
                                    buChar = pPr.find('a:buChar', ns)
                                    buAutoNum = pPr.find('a:buAutoNum', ns)
                                    if buChar is not None or buAutoNum is not None:
                                        bullet_count += 1

                        if total_paras > 0 and bullet_count >= 2:
                            has_bullets = True
                            print(f"PASS: Component 3b -- {bullet_count}/{total_paras} paragraphs have bullets (0.20 pts)")
                            comp3_score += 0.20
                        else:
                            print(f"FAIL: Component 3b -- Only {bullet_count}/{total_paras} paragraphs have bullets")
            except Exception as e:
                print(f"ERROR: Component 3b XML check -- {e}")
                # Fallback: check if text has bullet-like characters
                for p in non_empty:
                    txt = p.text.strip()
                    if txt.startswith(('- ', '* ', chr(8226), chr(9679))):
                        has_bullets = True
                        break
                if has_bullets:
                    print(f"PASS: Component 3b -- Bullet characters found in text (fallback) (0.20 pts)")
                    comp3_score += 0.20
                else:
                    print(f"FAIL: Component 3b -- No bullet formatting detected")
        else:
            print(f"FAIL: Component 3 -- No right-side text box with multiple items found")

        total_score += comp3_score
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
