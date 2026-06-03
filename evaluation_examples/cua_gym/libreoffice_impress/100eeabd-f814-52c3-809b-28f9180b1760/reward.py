"""
Reward Script: Change bar chart to column chart with blue gradient colors
Task ID: impress_tct_039
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Chart type changed from horizontal bar to vertical column
  Component 2 (0.4): Data points use blue color scheme (shades of blue)
  Component 3 (0.2): Data integrity preserved (categories + values unchanged)
"""

import os
import zipfile
import re

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_039'


def is_blue_shade(hex_color):
    """Check if a hex color is a shade of blue.
    Blue shades have higher blue channel relative to red and green.
    """
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        # Blue channel must be dominant or at least significantly present
        # For blue gradient: B should be > R and B should be > G
        return b > r and b > g
    except (ValueError, IndexError):
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be a valid pptx
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: slide 3 must exist and have a chart
    if len(prs.slides) < 3:
        print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        print("CRITICAL: No chart found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    chart = chart_shape.chart

    # =========================================================
    # Component 1: Chart type is column (vertical bars) (0.4 pts)
    # Initial: BAR_CLUSTERED (57) with barDir="bar" (horizontal)
    # Golden: COLUMN_CLUSTERED (51) with barDir="col" (vertical)
    # =========================================================
    try:
        # Check via XML for barDir value - most reliable
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_xml_files = [f for f in zf.namelist() if f.startswith('ppt/charts/chart') and f.endswith('.xml')]
            if not chart_xml_files:
                print("FAIL: Component 1 -- No chart XML found in archive")
            else:
                chart_content = zf.read(chart_xml_files[0]).decode()
                bar_dirs = re.findall(r'<c:barDir val="([^"]+)"', chart_content)
                if bar_dirs and bar_dirs[0] == 'col':
                    print(f"PASS: Component 1 -- Chart direction is 'col' (vertical column) (0.4 pts)")
                    total_score += 0.4
                else:
                    actual_dir = bar_dirs[0] if bar_dirs else 'unknown'
                    print(f"FAIL: Component 1 -- Expected barDir='col', found '{actual_dir}'")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # =========================================================
    # Component 2: Data points use blue color scheme (0.4 pts)
    # Initial: single solid red fill CC0000
    # Golden: 5 per-point blue shades via dPt elements
    # =========================================================
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_content = zf.read(chart_xml_files[0]).decode()

            # Find all per-data-point color definitions
            dpts = re.findall(r'<c:dPt>.*?</c:dPt>', chart_content, re.DOTALL)
            # Also check series-level solid fill colors
            all_colors = re.findall(r'<a:srgbClr val="([^"]+)"', chart_content)

            if not all_colors:
                print("FAIL: Component 2 -- No colors found in chart XML")
            else:
                blue_count = sum(1 for c in all_colors if is_blue_shade(c))
                non_blue = [c for c in all_colors if not is_blue_shade(c)]

                if blue_count == 0:
                    print(f"FAIL: Component 2 -- No blue colors found. Colors: {all_colors}")
                elif blue_count == len(all_colors) and blue_count >= 1:
                    # All colors are blue shades
                    # Check for gradient-like variety (multiple distinct blues)
                    unique_blues = set(c for c in all_colors if is_blue_shade(c))
                    if len(unique_blues) >= 3:
                        print(f"PASS: Component 2 -- Blue gradient scheme with {len(unique_blues)} distinct shades: {sorted(unique_blues)} (0.4 pts)")
                        total_score += 0.4
                    elif len(unique_blues) >= 1:
                        # At least blue, but not enough variety for "gradient"
                        print(f"PARTIAL: Component 2 -- Blue colors found but only {len(unique_blues)} distinct shades (need 3+ for gradient) (0.2 pts)")
                        total_score += 0.2
                else:
                    # Some blue, some not
                    if blue_count >= 3:
                        print(f"PARTIAL: Component 2 -- {blue_count}/{len(all_colors)} colors are blue. Non-blue: {non_blue} (0.2 pts)")
                        total_score += 0.2
                    else:
                        print(f"FAIL: Component 2 -- Only {blue_count}/{len(all_colors)} colors are blue. All colors: {all_colors}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # =========================================================
    # Component 3: Chart is column AND data integrity preserved (0.2 pts)
    # This is a compound check: the chart must already be a column chart
    # (task-introduced change) AND the original data must be intact.
    # This ensures we don't award points for data that was already there.
    # =========================================================
    try:
        expected_categories = ['Premium Electronics', 'Organic Food & Beverages', 'Home Automation', 'Fitness Equipment', 'Sustainable Fashion']
        expected_values = [4520.0, 3870.0, 3150.0, 2680.0, 2340.0]

        plot = chart.plots[0]
        actual_categories = [str(c) for c in plot.categories]
        actual_values = list(chart.series[0].values)

        cats_match = actual_categories == expected_categories
        vals_match = actual_values == expected_values

        # Re-check barDir to anchor this component to the task change
        with zipfile.ZipFile(file_path, 'r') as zf:
            chart_xml_files_c3 = [f for f in zf.namelist() if f.startswith('ppt/charts/chart') and f.endswith('.xml')]
            chart_content_c3 = zf.read(chart_xml_files_c3[0]).decode()
            bar_dirs_c3 = re.findall(r'<c:barDir val="([^"]+)"', chart_content_c3)
            is_column = bar_dirs_c3 and bar_dirs_c3[0] == 'col'

        if is_column and cats_match and vals_match:
            print(f"PASS: Component 3 -- Column chart with data integrity preserved: {len(actual_categories)} categories, values match (0.2 pts)")
            total_score += 0.2
        else:
            if not is_column:
                print(f"FAIL: Component 3 -- Chart is not column type, skipping data integrity check")
            if not cats_match:
                print(f"FAIL: Component 3 -- Categories mismatch. Expected: {expected_categories}, Got: {actual_categories}")
            if not vals_match:
                print(f"FAIL: Component 3 -- Values mismatch. Expected: {expected_values}, Got: {actual_values}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
