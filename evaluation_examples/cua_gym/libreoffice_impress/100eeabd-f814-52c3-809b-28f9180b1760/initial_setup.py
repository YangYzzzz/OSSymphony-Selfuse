"""
Initial Setup: Sales Trends presentation with horizontal bar chart on slide 3
Task ID: impress_tct_039
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_039'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Sales Trends Q4 2025"
    slide1.placeholders[1].text = "Regional Product Performance Review\nPrepared by Analytics Team"

    # --- Slide 2: Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Market Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    bullets = [
        "Total revenue grew 12.4% year-over-year, reaching $48.7M in Q4 2025",
        "Five core product lines drove 87% of total revenue",
        "Premium Electronics segment showed the strongest growth at 18.2%",
        "Organic Food & Beverages maintained steady demand across all regions",
        "Home Automation products gained significant traction in urban markets",
    ]
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = text
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 3: Horizontal Bar Chart ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title for chart slide
    title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11), Inches(0.8))
    tf3 = title_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Product Sales Performance (Units Sold)"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # Create horizontal bar chart
    chart_data = CategoryChartData()
    chart_data.categories = [
        'Premium Electronics',
        'Organic Food & Beverages',
        'Home Automation',
        'Fitness Equipment',
        'Sustainable Fashion',
    ]
    chart_data.add_series('Q4 2025 Sales', (4520, 3870, 3150, 2680, 2340))

    chart_frame = slide3.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,  # horizontal bar
        Inches(1.0), Inches(1.3),
        Inches(11.0), Inches(5.8),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the bars red
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = RGBColor(0xCC, 0x00, 0x00)

    # Category axis
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    category_axis.tick_labels.font.size = Pt(12)

    # Value axis
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.font.size = Pt(10)

    # --- Slide 4: Regional Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    title_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    tf4 = title_box4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Regional Sales Breakdown"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.size = Pt(32)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    # Add a table with regional data
    rows, cols = 6, 4
    tbl_shape = slide4.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8), Inches(10.5), Inches(4.0))
    table = tbl_shape.table

    headers = ['Region', 'Q4 Revenue ($M)', 'Growth (%)', 'Market Share (%)']
    data_rows = [
        ['North America', '$18.2M', '+14.1%', '37.4%'],
        ['Europe', '$12.8M', '+9.7%', '26.3%'],
        ['Asia Pacific', '$10.5M', '+18.6%', '21.6%'],
        ['Latin America', '$4.3M', '+11.2%', '8.8%'],
        ['Middle East & Africa', '$2.9M', '+7.3%', '5.9%'],
    ]

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(12)

    # --- Slide 5: Key Takeaways ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    title_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
    tf5 = title_box5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Key Takeaways & Next Steps"
    p5.alignment = PP_ALIGN.LEFT
    run5 = p5.runs[0]
    run5.font.size = Pt(32)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0x2E, 0x3B, 0x4E)

    content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
    tf5c = content_box5.text_frame
    tf5c.word_wrap = True
    takeaways = [
        "Premium Electronics and Home Automation are our fastest-growing segments",
        "Asia Pacific region shows the highest growth potential at 18.6% YoY",
        "Invest in Sustainable Fashion marketing to capture growing eco-conscious demographic",
        "Expand Fitness Equipment distribution channels in Latin America",
        "Target 15% overall revenue growth for Q1 2026 with enhanced digital presence",
    ]
    for i, text in enumerate(takeaways):
        if i == 0:
            p = tf5c.paragraphs[0]
        else:
            p = tf5c.add_paragraph()
        p.text = text
        p.space_after = Pt(10)
        for r in p.runs:
            r.font.size = Pt(18)
            r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
