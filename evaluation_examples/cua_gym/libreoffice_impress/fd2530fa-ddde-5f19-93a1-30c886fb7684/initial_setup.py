"""
Initial Setup: Infeasible task - attempt to password-protect Confidential_Slides.pptx
Task ID: osworld_multi_apps_impress_infeasible_009
Domain: libreoffice_impress

Creates a realistic confidential business presentation that is open in LibreOffice Impress.
The task (applying password protection via LibreOffice Impress built-in UI for PPTX) is
infeasible because LibreOffice Impress does not support open-password protection for PPTX format.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_multi_apps_impress_infeasible_009'
# The task instruction references 'Confidential_Slides.pptx' by name
OUTPUT = f'{WORKDIR}/Confidential_Slides.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Confidential Business Strategy 2025"
    slide1.placeholders[1].text = "Q2 Executive Review\nFor Internal Use Only"

    # Title formatting
    title_tf = slide1.shapes.title.text_frame
    for para in title_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(36)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7D)

    # Subtitle formatting
    sub_tf = slide1.placeholders[1].text_frame
    for para in sub_tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    # Background
    bg1 = slide1.background.fill
    bg1.solid()
    bg1.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)

    # ---- Slide 2: Financial Overview ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Q2 Financial Overview"

    content_tf = slide2.placeholders[1].text_frame
    content_tf.text = "Revenue: $4.2M (up 18% YoY)"

    lines = [
        "Operating Expenses: $2.8M",
        "Net Profit Margin: 33%",
        "Customer Acquisition Cost: $142 per user",
        "Monthly Recurring Revenue: $350K",
        "Projected Q3 Revenue: $5.1M",
    ]
    for line in lines:
        para = content_tf.add_paragraph()
        para.text = line
        para.level = 1
        for run in para.runs:
            run.font.size = Pt(18)

    title2_tf = slide2.shapes.title.text_frame
    for para in title2_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7D)

    # ---- Slide 3: Key Initiatives ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Strategic Initiatives"

    init_tf = slide3.placeholders[1].text_frame
    init_tf.text = "Product Development"
    initiatives = [
        "Launch of AI-powered analytics dashboard - Q3 2025",
        "Mobile app redesign targeting enterprise clients",
        "Expansion into APAC market (Singapore, Tokyo)",
        "Partnership with TechVision Corp for cloud integration",
        "Hiring 25 additional engineers by September 2025",
    ]
    for item in initiatives:
        para = init_tf.add_paragraph()
        para.text = item
        para.level = 1
        for run in para.runs:
            run.font.size = Pt(17)

    title3_tf = slide3.shapes.title.text_frame
    for para in title3_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7D)

    # ---- Slide 4: Risk Assessment ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Risk Assessment & Mitigation"

    risk_tf = slide4.placeholders[1].text_frame
    risk_tf.text = "Identified Risks (Confidential)"
    risks = [
        "Regulatory compliance in EU market - Mitigation: Legal review underway",
        "Supply chain disruption - Mitigation: Dual-supplier strategy active",
        "Competitor product launch (EstimatedQ4) - Mitigation: Accelerate roadmap",
        "Key personnel retention - Mitigation: Revised compensation packages",
        "Data security vulnerabilities - Mitigation: Third-party audit scheduled",
    ]
    for r in risks:
        para = risk_tf.add_paragraph()
        para.text = r
        para.level = 1
        for run in para.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    title4_tf = slide4.shapes.title.text_frame
    for para in title4_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7D)

    # ---- Slide 5: Org & Personnel ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Executive Team & Reporting"

    org_tf = slide5.placeholders[1].text_frame
    org_tf.text = "Leadership Roster"
    personnel = [
        "CEO: Alexandra Reeves - a.reeves@corp-internal.com",
        "CFO: Jonathan Mercer - j.mercer@corp-internal.com",
        "CTO: Priya Nair - p.nair@corp-internal.com",
        "VP Sales: Daniel Okafor - d.okafor@corp-internal.com",
        "VP Engineering: Mei-Ling Zhou - m.zhou@corp-internal.com",
    ]
    for p_line in personnel:
        para = org_tf.add_paragraph()
        para.text = p_line
        para.level = 1
        for run in para.runs:
            run.font.size = Pt(17)

    title5_tf = slide5.shapes.title.text_frame
    for para in title5_tf.paragraphs:
        for run in para.runs:
            run.font.bold = True
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x7D)

    # Add speaker notes to first slide
    slide1.notes_slide.notes_text_frame.text = (
        "CONFIDENTIAL - Do not distribute outside of executive team.\n"
        "Prepared by Strategy Office, March 2025."
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
