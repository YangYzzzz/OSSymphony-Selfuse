"""
Reward Script: Event planning presentation with 9 slides
Task ID: impress_wf_038
Domain: libreoffice_impress
Scoring:
  C1: File exists on Desktop (0.10)
  C2: Exactly 9 slides (0.10)
  C3: Slide 1 title "Annual Gala 2024" (0.10)
  C4: Slide 3 has budget table (0.10)
  C5: Slide 4 has colored Gantt bars (>=5 rectangles with fill) (0.15)
  C6: Slide 5 has venue layout with mixed shapes (0.15)
  C7: Slide 6 has vendor contacts table (0.10)
  C8: Slide 7 has 3-column menu layout (0.10)
  C9: Slide 9 has table + colored circle status indicators (0.10)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_038'
FILE_PATH = os.path.join(WORKDIR, 'Desktop', 'Gala_Planning.pptx')


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and be loadable
    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File exists on Desktop with correct name (0.10 points)
    # The file existence is already confirmed above, but this component scores
    # that the file is specifically at Desktop/Gala_Planning.pptx (not elsewhere)
    # This FAILS on initial_env because no such file exists on Desktop
    try:
        if os.path.exists(file_path) and num_slides >= 1:
            print(f"PASS: Component 1 — Gala_Planning.pptx exists on Desktop with {num_slides} slides (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 1 — File not found or empty")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Exactly 9 slides (0.10 points)
    try:
        if num_slides == 9:
            print(f"PASS: Component 2 — Presentation has exactly 9 slides (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 2 — Expected 9 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 1 title contains "Annual Gala 2024" (0.10 points)
    try:
        slide1 = prs.slides[0]
        slide1_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    slide1_text += para.text + " "
        slide1_text = slide1_text.strip()
        if "annual gala 2024" in slide1_text.lower():
            print(f"PASS: Component 3 — Slide 1 contains 'Annual Gala 2024' (0.10 pts)")
            total_score += 0.10
        else:
            print(f"FAIL: Component 3 — Slide 1 text: '{slide1_text[:100]}', expected 'Annual Gala 2024'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 3 has a budget breakdown table (0.10 points)
    # Must have a TABLE shape with budget-related headers
    try:
        if num_slides >= 3:
            slide3 = prs.slides[2]
            table_found = False
            for shape in slide3.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    # Check it has at least a few rows and columns
                    if len(table.rows) >= 3 and len(table.columns) >= 2:
                        # Check header row contains budget-related terms
                        headers = [table.cell(0, c).text.lower() for c in range(len(table.columns))]
                        header_text = " ".join(headers)
                        if "category" in header_text or "cost" in header_text or "budget" in header_text:
                            table_found = True
                            print(f"PASS: Component 4 — Slide 3 has budget table ({len(table.rows)}x{len(table.columns)}) (0.10 pts)")
                            total_score += 0.10
                            break
            if not table_found:
                print(f"FAIL: Component 4 — No budget table found on slide 3")
        else:
            print(f"FAIL: Component 4 — Not enough slides (need >= 3)")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 4 has colored Gantt bars (>=5 colored rectangles) (0.15 points)
    try:
        if num_slides >= 4:
            slide4 = prs.slides[3]
            colored_rects = 0
            for shape in slide4.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "Rectangle" in shape.name:
                    try:
                        fill = shape.fill
                        if fill.type is not None and fill.type == 1:  # SOLID fill
                            colored_rects += 1
                    except Exception:
                        pass
            if colored_rects >= 5:
                print(f"PASS: Component 5 — Slide 4 has {colored_rects} colored Gantt bars (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 — Expected >= 5 colored rectangles on slide 4, found {colored_rects}")
        else:
            print(f"FAIL: Component 5 — Not enough slides (need >= 4)")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: Slide 5 has venue layout with mixed shapes (rectangles, ovals, lines) (0.15 points)
    try:
        if num_slides >= 5:
            slide5 = prs.slides[4]
            has_ovals = False
            has_rects = False
            has_lines = False
            for shape in slide5.shapes:
                st = shape.shape_type
                if st == MSO_SHAPE_TYPE.AUTO_SHAPE and "Oval" in shape.name:
                    has_ovals = True
                elif st == MSO_SHAPE_TYPE.AUTO_SHAPE and "Rectangle" in shape.name:
                    has_rects = True
                elif st == MSO_SHAPE_TYPE.LINE:
                    has_lines = True

            if has_ovals and has_rects and has_lines:
                print(f"PASS: Component 6 — Slide 5 has ovals, rectangles, and lines (venue layout) (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 6 — Slide 5 missing shapes: ovals={has_ovals}, rects={has_rects}, lines={has_lines}")
        else:
            print(f"FAIL: Component 6 — Not enough slides (need >= 5)")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    # Component 7: Slide 6 has a vendor contacts table (0.10 points)
    try:
        if num_slides >= 6:
            slide6 = prs.slides[5]
            vendor_table_found = False
            for shape in slide6.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    if len(table.rows) >= 3 and len(table.columns) >= 3:
                        headers = [table.cell(0, c).text.lower() for c in range(len(table.columns))]
                        header_text = " ".join(headers)
                        if "vendor" in header_text or "contact" in header_text:
                            vendor_table_found = True
                            print(f"PASS: Component 7 — Slide 6 has vendor contacts table ({len(table.rows)}x{len(table.columns)}) (0.10 pts)")
                            total_score += 0.10
                            break
            if not vendor_table_found:
                print(f"FAIL: Component 7 — No vendor contacts table found on slide 6")
        else:
            print(f"FAIL: Component 7 — Not enough slides (need >= 6)")
    except Exception as e:
        print(f"ERROR: Component 7 — {e}")

    # Component 8: Slide 7 has 3-column menu layout (Appetizers, Mains, Desserts) (0.10 points)
    try:
        if num_slides >= 7:
            slide7 = prs.slides[6]
            all_text = ""
            text_boxes = []
            for shape in slide7.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip().lower()
                    all_text += " " + txt
                    text_boxes.append(txt)

            has_appetizers = "appetizer" in all_text
            has_mains = "main" in all_text
            has_desserts = "dessert" in all_text

            if has_appetizers and has_mains and has_desserts:
                print(f"PASS: Component 8 — Slide 7 has 3-column menu (Appetizers, Mains, Desserts) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 8 — Missing menu columns: appetizers={has_appetizers}, mains={has_mains}, desserts={has_desserts}")
        else:
            print(f"FAIL: Component 8 — Not enough slides (need >= 7)")
    except Exception as e:
        print(f"ERROR: Component 8 — {e}")

    # Component 9: Slide 9 has table + colored circle status indicators (green/yellow/red) (0.10 points)
    try:
        if num_slides >= 9:
            slide9 = prs.slides[8]
            has_table = False
            green_count = 0
            yellow_count = 0
            red_count = 0

            for shape in slide9.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    if len(table.rows) >= 3:
                        has_table = True

                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and "Oval" in shape.name:
                    try:
                        fill = shape.fill
                        if fill.type is not None and fill.type == 1:
                            rgb = str(fill.fore_color.rgb).upper()
                            r_val = int(rgb[0:2], 16)
                            g_val = int(rgb[2:4], 16)
                            b_val = int(rgb[4:6], 16)
                            # Classify: green-ish, yellow-ish, red-ish
                            if g_val > 150 and r_val < 100 and b_val < 100:
                                green_count += 1
                            elif r_val > 150 and g_val > 150 and b_val < 100:
                                yellow_count += 1
                            elif r_val > 150 and g_val < 100 and b_val < 100:
                                red_count += 1
                    except Exception:
                        pass

            if has_table and green_count >= 1 and yellow_count >= 1 and red_count >= 1:
                print(f"PASS: Component 9 — Slide 9 has table + colored status circles (G={green_count}, Y={yellow_count}, R={red_count}) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 9 — table={has_table}, green={green_count}, yellow={yellow_count}, red={red_count}")
        else:
            print(f"FAIL: Component 9 — Not enough slides (need >= 9)")
    except Exception as e:
        print(f"ERROR: Component 9 — {e}")

    final_score = min(total_score, 1.0)
    # Round to avoid floating point drift
    final_score = round(final_score, 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
if not os.path.exists(FILE_PATH):
    print(f"File not found: {FILE_PATH}")
    print("REWARD: 0.0")
else:
    verify_task(FILE_PATH)
