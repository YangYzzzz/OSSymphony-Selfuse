"""
Reward Script: Add a team members credits slide with a name/contribution table
Task ID: impress_stu_035
Domain: libreoffice_impress
Scoring:
  C1 (0.15): Slide count is 11 (was 10)
  C2 (0.20): Slide 2 contains "Team Members" title text
  C3 (0.15): Slide 2 has a table with 2 columns headed "Name" and "Contribution"
  C4 (0.10): Table has exactly 5 rows (1 header + 4 data)
  C5 (0.20): All 4 team member names present in correct rows
  C6 (0.20): All 4 contribution descriptions present in correct rows
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_035'

# Expected table data (row index 1-4, after header row 0)
EXPECTED_MEMBERS = [
    ("Alex Kim", "Research & Literature Review"),
    ("Maria Santos", "Data Analysis & Charts"),
    ("James Okafor", "Slide Design & Formatting"),
    ("Sarah Liu", "Presentation & Editing"),
]


def find_table_on_slide(slide):
    """Find the first table shape on the slide, return table or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_all_text(slide):
    """Get concatenated text from all text-bearing shapes on the slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text.strip())
    return texts


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide count is 11 (was 10 initially) — 0.15 pts
    try:
        if num_slides == 11:
            print(f"PASS: Component 1 — slide count is 11 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — expected 11 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Need at least 2 slides to check slide 2
    if num_slides < 2:
        print("CRITICAL: Not enough slides to check slide 2")
        print(f"REWARD: {total_score}")
        return total_score

    slide2 = prs.slides[1]  # 0-indexed, slide 2

    # Component 2: Slide 2 has "Team Members" title text — 0.20 pts
    try:
        all_texts = get_all_text(slide2)
        has_team_members = any("team members" in t.lower() for t in all_texts)
        if has_team_members:
            print(f"PASS: Component 2 — 'Team Members' found on slide 2 (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 — 'Team Members' not found on slide 2. Texts: {all_texts}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 2 has a table with 2 columns headed "Name" and "Contribution" — 0.15 pts
    table = None
    try:
        table = find_table_on_slide(slide2)
        if table is not None:
            ncols = len(table.columns)
            if ncols == 2:
                h0 = table.cell(0, 0).text.strip().lower()
                h1 = table.cell(0, 1).text.strip().lower()
                if h0 == "name" and h1 == "contribution":
                    print(f"PASS: Component 3 — table with 'Name'/'Contribution' headers (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 3 — headers are '{h0}'/'{h1}', expected 'name'/'contribution'")
            else:
                print(f"FAIL: Component 3 — table has {ncols} columns, expected 2")
        else:
            print(f"FAIL: Component 3 — no table found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Table has 5 rows (header + 4 data) — 0.10 pts
    try:
        if table is not None:
            nrows = len(table.rows)
            if nrows == 5:
                print(f"PASS: Component 4 — table has 5 rows (1 header + 4 data) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 4 — table has {nrows} rows, expected 5")
        else:
            print(f"FAIL: Component 4 — no table to check row count")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: All 4 team member names in correct rows — 0.20 pts
    try:
        if table is not None and len(table.rows) >= 5:
            names_correct = 0
            for i, (expected_name, _) in enumerate(EXPECTED_MEMBERS):
                actual_name = table.cell(i + 1, 0).text.strip()
                if actual_name.lower() == expected_name.lower():
                    names_correct += 1
                else:
                    print(f"  Name row {i+1}: expected '{expected_name}', found '{actual_name}'")
            if names_correct == 4:
                print(f"PASS: Component 5 — all 4 names correct (0.20 pts)")
                total_score += 0.20
            elif names_correct > 0:
                # Partial credit: 0.05 per correct name
                partial = names_correct * 0.05
                if partial > 0:
                    print(f"PARTIAL: Component 5 — {names_correct}/4 names correct ({partial:.2f} pts)")
                    total_score += partial
            else:
                print(f"FAIL: Component 5 — 0/4 names correct")
        else:
            print(f"FAIL: Component 5 — no table or insufficient rows")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: All 4 contribution descriptions in correct rows — 0.20 pts
    try:
        if table is not None and len(table.rows) >= 5:
            contribs_correct = 0
            for i, (_, expected_contrib) in enumerate(EXPECTED_MEMBERS):
                actual_contrib = table.cell(i + 1, 1).text.strip()
                if actual_contrib.lower() == expected_contrib.lower():
                    contribs_correct += 1
                else:
                    print(f"  Contrib row {i+1}: expected '{expected_contrib}', found '{actual_contrib}'")
            if contribs_correct == 4:
                print(f"PASS: Component 6 — all 4 contributions correct (0.20 pts)")
                total_score += 0.20
            elif contribs_correct > 0:
                # Partial credit: 0.05 per correct contribution
                partial = contribs_correct * 0.05
                if partial > 0:
                    print(f"PARTIAL: Component 6 — {contribs_correct}/4 contributions correct ({partial:.2f} pts)")
                    total_score += partial
            else:
                print(f"FAIL: Component 6 — 0/4 contributions correct")
        else:
            print(f"FAIL: Component 6 — no table or insufficient rows")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
