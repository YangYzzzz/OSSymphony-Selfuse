"""
Initial Setup: Create a 10-slide Group Capstone presentation
Task ID: impress_stu_035
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_035'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body text on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.text = body_lines[0] if body_lines else ""
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line
            break


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Group Capstone Project"
    for ph in slide1.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Urban Sustainability: Analyzing Green Infrastructure Impact on City Livability"
            break

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Introduction", [
        "Rapid urbanization presents significant challenges for sustainable development",
        "Green infrastructure (GI) includes parks, green roofs, urban forests, and bioswales",
        "This study examines GI adoption across 12 metropolitan areas from 2018 to 2024",
        "Research question: How does green infrastructure investment correlate with livability indices?",
    ])

    # --- Slide 3: Literature Review ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Literature Review", [
        "Benedict & McMahon (2006): Foundational framework for GI planning",
        "Tzoulas et al. (2007): Health benefits of urban green spaces",
        "Demuzere et al. (2014): Ecosystem services provided by GI in cities",
        "Kabisch et al. (2016): Equity concerns in green space distribution",
        "Recent meta-analyses confirm positive correlation between GI and property values",
    ])

    # --- Slide 4: Methodology ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Methodology", [
        "Mixed-methods approach combining quantitative and qualitative analysis",
        "Data sources: EPA GI databases, Census Bureau, city planning documents",
        "Quantitative: Regression analysis on GI spending vs. livability metrics",
        "Qualitative: Semi-structured interviews with 24 urban planners",
        "Statistical tools: R (version 4.3) and Python for data processing",
    ])

    # --- Slide 5: Data Analysis ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Data Analysis", [
        "Analyzed 12 cities with populations ranging from 250K to 3.5M",
        "GI investment per capita ranged from $42 to $187 annually",
        "Livability index improvements correlated at r = 0.78 (p < 0.001)",
        "Air quality improvements showed strongest individual correlation (r = 0.84)",
        "Heat island mitigation averaged 2.3°F reduction in high-GI areas",
    ])

    # --- Slide 6: Results ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Results", [
        "Cities with >$100/capita GI spending showed 23% higher livability scores",
        "Green roof adoption reduced stormwater runoff by 34% in test neighborhoods",
        "Urban tree canopy above 25% correlated with 15% lower hospitalization rates",
        "Property values within 500m of new GI projects increased by average 8.2%",
        "Community engagement scores improved 41% in neighborhoods with GI projects",
    ])

    # --- Slide 7: Discussion ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Discussion", [
        "Results strongly support the economic case for green infrastructure investment",
        "Equity gaps persist: lower-income neighborhoods receive 37% less GI funding",
        "Maintenance costs remain a challenge for long-term sustainability",
        "Policy recommendations: dedicated GI funding, equity mandates, public-private partnerships",
        "Limitations: 6-year timeframe may not capture long-term ecological impacts",
    ])

    # --- Slide 8: Conclusion ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Conclusion", [
        "Green infrastructure delivers measurable improvements to urban livability",
        "Investment threshold of $80/capita appears necessary for significant impact",
        "Equitable distribution requires intentional policy intervention",
        "Future research should examine 10+ year longitudinal outcomes",
        "Collaboration between planners, ecologists, and communities is essential",
    ])

    # --- Slide 9: References ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide9, "References", [
        "Benedict, M. & McMahon, E. (2006). Green Infrastructure. Island Press.",
        "Tzoulas, K. et al. (2007). Urban Forestry & Urban Greening, 6(3), 167-177.",
        "Demuzere, M. et al. (2014). Environmental Science & Policy, 37, 55-67.",
        "Kabisch, N. et al. (2016). Landscape and Urban Planning, 156, 12-21.",
        "EPA (2023). Green Infrastructure Assessment Report. Washington, DC.",
    ])

    # --- Slide 10: Q&A ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide10, "Questions & Discussion", [
        "Thank you for your attention!",
        "",
        "We welcome your questions and feedback.",
        "",
        "Contact: capstone-team@university.edu",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
