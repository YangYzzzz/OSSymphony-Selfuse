"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 4, I need the footer to show the date exactly as “2024-01-01” (no automatic updates). What’s the quickest way to drop that fixed timestamp into LibreOffice Impress?
Generated: 2025-09-10 21:47:02
Status: success
Model: azure-o3
Total Steps: 2
"""

from pptx import Presentation
from lxml import etree
import os, re


def verify_fixed_date_footer(file_path: str, expected_date: str = "2024-01-01") -> float:
    """Verify that slide 4 contains the exact fixed date in the footer (not an automatic date field).

    Scoring (progressive):
        0.6 pts – Slide 4 contains a shape whose **entire text** equals the expected date string.
        0.4 pts – That shape’s XML contains **no** <a:fld> elements (i.e., it is fixed text, not an auto-date field).
        1.0 pts – Both conditions satisfied.
    """

    print(f"Verifying fixed date footer in: {file_path}\nExpected date: {expected_date}")

    max_score = 1.0
    score = 0.0

    # ---------- 1. Load presentation ----------
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ---------- 2. Ensure slide 4 exists ----------
    target_idx = 3  # zero-based index for slide 4
    if len(prs.slides) <= target_idx:
        print(f"✗ Presentation has only {len(prs.slides)} slides – slide 4 missing")
        return 0.0
    slide4 = prs.slides[target_idx]

    # ---------- 3. Search for exact date text on slide ----------
    date_regex = re.compile(re.escape(expected_date))
    candidate_shapes = []

    for shape in slide4.shapes:
        text = ""
        # Try .text attribute first; fall back to text_frame
        if hasattr(shape, "text"):
            text = shape.text.strip()
        if not text and getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()

        if text:
            print(f"  Found text in shape: '{text}'")
            if date_regex.fullmatch(text):
                candidate_shapes.append(shape)

    if not candidate_shapes:
        print("✗ No shape on slide 4 contains the exact expected date")
        return 0.0

    print(f"✓ Found {len(candidate_shapes)} shape(s) with the expected date text (0.6 points)")
    score += 0.6

    # ---------- 4. Confirm the date is fixed (not automatic) ----------
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    fixed_shape_found = False

    for shp in candidate_shapes:
        try:
            xml_root = etree.fromstring(shp._element.xml.encode())
            fld_elems = xml_root.xpath(".//a:fld", namespaces=ns)
            if not fld_elems:
                fixed_shape_found = True
                break
            else:
                print("  Shape contains automatic field – ignored for fixed-date requirement")
        except Exception as e:
            print(f"  Error inspecting shape XML: {e}")

    if fixed_shape_found:
        print("✓ Verified that the date text is fixed (not an automatic field) (0.4 points)")
        score += 0.4
    else:
        print("✗ All matching shapes appear to be automatic fields – fixed date not found")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_4_i_need_the_footer_to_show_the_date_exactly_as_2024_01_01_no_automatic_updates_whats_the_q_golden.pptx"
    reward = verify_fixed_date_footer(FILE_PATH)
    print(f"REWARD: {reward}")
