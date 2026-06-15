"""
Reward Script: Partner ecosystem diagram on slide 8
Task ID: impress_sales_080
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Central "Our Platform" circle with #2B6CB0 fill and white text
  Component 2 (0.30): 6 partner circles with correct names
  Component 3 (0.20): Partner circles styled with #F5F5F5 fill and #333333 text
  Component 4 (0.25): 6 connector lines from center to partners
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_080'

EXPECTED_PARTNERS = {'Salesforce', 'HubSpot', 'Slack', 'Jira', 'Zendesk', 'Shopify'}


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        import time
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify partner ecosystem slide on slide 8.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[7]  # slide 8 (0-indexed)

    # Collect shapes by type for analysis
    ovals = []
    connectors = []
    for shape in slide.shapes:
        # AUTO_SHAPE (type 1) - circles/ovals
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Get fill color
                fill_color = None
                try:
                    if shape.fill.type is not None and shape.fill.type == 1:  # SOLID
                        fill_color = str(shape.fill.fore_color.rgb)
                except Exception:
                    pass

                # Get text color from first non-empty run
                text_color = None
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            try:
                                if run.font.color.type is not None:
                                    text_color = str(run.font.color.rgb)
                            except Exception:
                                pass
                            break
                    if text_color is not None:
                        break

                ovals.append({
                    'text': text,
                    'fill_color': fill_color,
                    'text_color': text_color,
                    'width': shape.width,
                    'height': shape.height,
                })

        # LINE / CONNECTOR (type 9) - connector lines
        if shape.shape_type == 9:  # MSO_SHAPE_TYPE.LINE
            connectors.append(shape)

    print(f"INFO: Found {len(ovals)} labeled ovals, {len(connectors)} connector lines on slide 8")

    # ---------------------------------------------------------------
    # Component 1: Central "Our Platform" circle (0.25 points)
    # Must have #2B6CB0 fill and white (FFFFFF) text
    # ---------------------------------------------------------------
    try:
        center_circle = None
        for oval in ovals:
            if 'our platform' in oval['text'].lower():
                center_circle = oval
                break

        if center_circle is None:
            print("FAIL: Component 1 — No circle with 'Our Platform' text found on slide 8")
        else:
            comp1_score = 0.0
            # Check text exists (the circle was found with 'Our Platform' text)
            if center_circle['text'].lower().strip() == 'our platform':
                comp1_score += 0.05
                print(f"  PASS: 'Our Platform' text found")
            # Check fill color #2B6CB0
            if center_circle['fill_color'] and center_circle['fill_color'].upper() == '2B6CB0':
                comp1_score += 0.10
                print(f"  PASS: Center circle fill is #2B6CB0")
            else:
                print(f"  FAIL: Center circle fill is {center_circle['fill_color']}, expected #2B6CB0")
            # Check text color white (FFFFFF)
            if center_circle['text_color'] and center_circle['text_color'].upper() == 'FFFFFF':
                comp1_score += 0.10
                print(f"  PASS: Center circle text is white (#FFFFFF)")
            else:
                print(f"  FAIL: Center circle text color is {center_circle['text_color']}, expected #FFFFFF")

            if comp1_score > 0:
                total_score += comp1_score
            print(f"PASS: Component 1 — Central 'Our Platform' circle ({comp1_score} pts)")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # ---------------------------------------------------------------
    # Component 2: 6 partner circles with correct names (0.30 points)
    # Each partner name found = 0.05 points
    # ---------------------------------------------------------------
    try:
        found_partners = set()
        for oval in ovals:
            text = oval['text'].strip()
            if text in EXPECTED_PARTNERS:
                found_partners.add(text)

        partner_count = len(found_partners)
        comp2_score = partner_count * 0.05
        if partner_count > 0:
            total_score += comp2_score

        if partner_count == 6:
            print(f"PASS: Component 2 — All 6 partner circles found: {found_partners} ({comp2_score} pts)")
        else:
            missing = EXPECTED_PARTNERS - found_partners
            print(f"FAIL: Component 2 — Found {partner_count}/6 partners. Missing: {missing} ({comp2_score} pts)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # ---------------------------------------------------------------
    # Component 3: Partner circle styling (0.20 points)
    # Each partner with correct fill (#F5F5F5) AND text color (#333333) = ~0.033 pts
    # ---------------------------------------------------------------
    try:
        styled_count = 0
        for oval in ovals:
            text = oval['text'].strip()
            if text in EXPECTED_PARTNERS:
                fill_ok = oval['fill_color'] and oval['fill_color'].upper() == 'F5F5F5'
                text_ok = oval['text_color'] and oval['text_color'].upper() == '333333'
                if fill_ok and text_ok:
                    styled_count += 1
                else:
                    print(f"  WARN: Partner '{text}' style — fill={oval['fill_color']} (expect F5F5F5), text={oval['text_color']} (expect 333333)")

        # Award proportional points
        if styled_count > 0:
            comp3_score = round(styled_count * (0.20 / 6), 4)
            if comp3_score > 0:
                total_score += comp3_score
            print(f"PASS: Component 3 — {styled_count}/6 partners correctly styled ({comp3_score} pts)")
        else:
            print(f"FAIL: Component 3 — No partner circles have correct #F5F5F5 fill + #333333 text (0.0 pts)")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # ---------------------------------------------------------------
    # Component 4: 6 connector lines (0.25 points)
    # Each connector line = ~0.0417 pts
    # ---------------------------------------------------------------
    try:
        num_connectors = len(connectors)
        if num_connectors >= 6:
            comp4_score = 0.25
            print(f"PASS: Component 4 — {num_connectors} connector lines found (>= 6 required) ({comp4_score} pts)")
        elif num_connectors > 0:
            comp4_score = round(num_connectors * (0.25 / 6), 4)
            print(f"PARTIAL: Component 4 — {num_connectors}/6 connector lines ({comp4_score} pts)")
        else:
            comp4_score = 0.0
            print(f"FAIL: Component 4 — No connector lines found on slide 8 (0.0 pts)")
        if comp4_score > 0:
            total_score += comp4_score
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
