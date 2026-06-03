"""
Initial Setup: Partner ecosystem pitch presentation with empty slide 8
Task ID: impress_sales_080
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_080'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = bullet_points[0]
    for bp in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = bp
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with title only (layout 5=blank, add title manually)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def create_initial():
    prs = Presentation()

    # Slide 1: Title
    add_title_slide(prs, "Ecosystem Pitch 2026",
                    "Connecting the Digital Workplace")

    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Market opportunity valued at $4.2B globally",
        "Platform adoption grew 180% year-over-year",
        "Strategic partnerships driving 65% of new revenue",
        "Customer retention rate at 94.7%",
    ])

    # Slide 3: Market Landscape
    add_content_slide(prs, "Market Landscape", [
        "Digital transformation accelerating across all industries",
        "SMBs increasingly adopting integrated SaaS platforms",
        "API-first architecture becoming the standard",
        "Integration demand growing at 23% CAGR",
    ])

    # Slide 4: Product Overview
    add_content_slide(prs, "Product Overview", [
        "Unified dashboard for cross-platform management",
        "Real-time data synchronization across 50+ tools",
        "Enterprise-grade security and compliance (SOC 2, GDPR)",
        "Custom workflow automation engine",
    ])

    # Slide 5: Customer Success Stories
    add_content_slide(prs, "Customer Success Stories", [
        "TechVenture Inc. reduced manual data entry by 73%",
        "GlobalRetail Corp saved $1.2M annually on integration costs",
        "HealthFirst achieved HIPAA compliance in 6 weeks",
        "EduConnect onboarded 15,000 users in a single quarter",
    ])

    # Slide 6: Revenue Growth
    add_content_slide(prs, "Revenue Growth", [
        "Q1 2025: $8.4M ARR (+42% YoY)",
        "Q2 2025: $10.1M ARR (+38% YoY)",
        "Q3 2025: $12.7M ARR (+51% YoY)",
        "Q4 2025: $15.3M ARR (+47% YoY)",
        "2026 Target: $24M ARR",
    ])

    # Slide 7: Technology Architecture
    add_content_slide(prs, "Technology Architecture", [
        "Microservices-based backend on Kubernetes",
        "GraphQL API gateway with rate limiting",
        "Event-driven architecture using Apache Kafka",
        "Multi-tenant data isolation with PostgreSQL",
    ])

    # Slide 8: Integration Ecosystem — TITLE ONLY, empty content
    # This is the slide the agent must complete
    add_title_only_slide(prs, "Integration Ecosystem")

    # Slide 9: Go-to-Market Strategy
    add_content_slide(prs, "Go-to-Market Strategy", [
        "Partner-led distribution through certified consultants",
        "Self-serve free tier with usage-based upgrades",
        "Enterprise sales team targeting Fortune 500",
        "Regional expansion into EMEA and APAC markets",
    ])

    # Slide 10: Next Steps & Timeline
    add_content_slide(prs, "Next Steps & Timeline", [
        "Q1 2026: Launch partner certification program",
        "Q2 2026: Release mobile SDK and companion app",
        "Q3 2026: Expand to 100+ native integrations",
        "Q4 2026: Series C fundraise target of $50M",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
