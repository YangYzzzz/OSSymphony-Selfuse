"""
Reward Script: Insert a chart on slide 2 using data from Sales_Data.xlsx A1:C7
Task ID: impress_tct_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Chart exists on slide 2
  Component 2 (0.30) - Chart has exactly 2 data series with correct values
  Component 3 (0.20) - Chart categories match the 6 month names
  Component 4 (0.20) - Slide count unchanged (5) and slide 2 title preserved
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_047'

# Expected data from Sales_Data.xlsx A1:C7
EXPECTED_SERIES_A = [12500.0, 14200.0, 11800.0, 16700.0, 15300.0, 18100.0]
EXPECTED_SERIES_B = [9800.0, 10500.0, 12300.0, 11200.0, 13800.0, 14600.0]
EXPECTED_CATEGORIES = ['January', 'February', 'March', 'April', 'May', 'June']


def persist_app_state(domain):
    """Try to save any unsaved LibreOffice state via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
    except ImportError:
        print("CRITICAL: python-pptx not available")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"PRECONDITION FAIL: Need at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed

    # Find chart shape(s) on slide 2
    chart_shapes = []
    for shape in slide2.shapes:
        if hasattr(shape, 'has_chart') and shape.has_chart:
            chart_shapes.append(shape)

    # Component 1: Chart exists on slide 2 (0.30 points)
    try:
        if len(chart_shapes) >= 1:
            print(f"PASS: Component 1 -- Chart found on slide 2 ({len(chart_shapes)} chart(s)) (0.30 pts)")
            total_score += 0.30
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 2")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If no chart, remaining checks cannot pass
    if len(chart_shapes) == 0:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shapes[0].chart

    # Component 2: Chart has 2 data series with correct values (0.30 points)
    try:
        series_count = len(chart.series)
        if series_count == 2:
            series_0_vals = list(chart.series[0].values)
            series_1_vals = list(chart.series[1].values)

            # Check if the values match expected (order may vary)
            match_a_0 = (series_0_vals == EXPECTED_SERIES_A)
            match_b_1 = (series_1_vals == EXPECTED_SERIES_B)
            match_a_1 = (series_1_vals == EXPECTED_SERIES_A)
            match_b_0 = (series_0_vals == EXPECTED_SERIES_B)

            if (match_a_0 and match_b_1) or (match_a_1 and match_b_0):
                print(f"PASS: Component 2 -- 2 series with correct data values (0.30 pts)")
                total_score += 0.30
            else:
                # Partial: correct series count but wrong values
                print(f"FAIL: Component 2 -- Series values don't match expected data")
                print(f"  Series 0: {series_0_vals}")
                print(f"  Series 1: {series_1_vals}")
        else:
            print(f"FAIL: Component 2 -- Expected 2 series, found {series_count}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart categories match month names (0.20 points)
    try:
        categories = [str(c) for c in chart.plots[0].categories]
        if categories == EXPECTED_CATEGORIES:
            print(f"PASS: Component 3 -- Categories match expected months (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 -- Expected categories {EXPECTED_CATEGORIES}, found {categories}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Slide count is 5 AND slide 2 still has "Sales Overview" text (0.20 points)
    # This is a compound check anchored to the chart insertion --
    # the chart must be added WITHOUT breaking the slide structure
    try:
        slide_count = len(prs.slides)
        # Check that "Sales Overview" text exists on slide 2
        sales_overview_count = sum(
            1 for shape in slide2.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            if 'Sales Overview' in para.text
        )
        has_sales_overview = sales_overview_count > 0

        has_chart = len(chart_shapes) >= 1
        if slide_count == 5 and has_sales_overview and has_chart:
            print(f"PASS: Component 4 -- Slide count=5, 'Sales Overview' preserved, chart present (0.20 pts)")
            total_score += 0.20
        else:
            reasons = []
            if slide_count != 5:
                reasons.append(f"slide count={slide_count} (expected 5)")
            if not has_sales_overview:
                reasons.append("'Sales Overview' text missing from slide 2")
            if not has_chart:
                reasons.append("no chart on slide 2")
            print(f"FAIL: Component 4 -- {'; '.join(reasons)}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
