"""
Initial Setup: Insert a chart on slide 2 using data from Sales_Data.xlsx
Task ID: impress_tct_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_047'
OUTPUT_PPTX = f'{WORKDIR}/{TASK_ID}.pptx'
OUTPUT_XLSX = f'{WORKDIR}/Sales_Data.xlsx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_sales_data():
    """Create Sales_Data.xlsx with month/product data in A1:C7."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = 'Sales'

    # Headers (row 1)
    ws['A1'] = 'Month'
    ws['B1'] = 'Product A'
    ws['C1'] = 'Product B'

    # Data rows (rows 2-7): January through June
    data = [
        ('January',   12500, 9800),
        ('February',  14200, 10500),
        ('March',     11800, 12300),
        ('April',     16700, 11200),
        ('May',       15300, 13800),
        ('June',      18100, 14600),
    ]
    for r, (month, prod_a, prod_b) in enumerate(data, 2):
        ws.cell(row=r, column=1, value=month)
        ws.cell(row=r, column=2, value=prod_a)
        ws.cell(row=r, column=3, value=prod_b)

    # Adjust column widths for readability
    ws.column_dimensions['A'].width = 12
    ws.column_dimensions['B'].width = 12
    ws.column_dimensions['C'].width = 12

    wb.save(OUTPUT_XLSX)
    print(f'Sales data created: {OUTPUT_XLSX}')


def create_presentation():
    """Create a 5-slide presentation. Slide 2 has title only, no chart."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Monthly Business Review"
    slide1.placeholders[1].text = "Q2 2025 Performance Summary"

    # --- Slide 2: Sales Overview (title only, NO chart) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Sales Overview"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    # --- Slide 3: Marketing Update ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf3t = txBox3_title.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Marketing Update"
    r3t = p3t.runs[0]
    r3t.font.name = "Arial"
    r3t.font.size = Pt(32)
    r3t.font.bold = True
    r3t.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    txBox3_body = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(5))
    tf3b = txBox3_body.text_frame
    tf3b.word_wrap = True
    tf3b.paragraphs[0].text = "Digital campaigns drove a 23% increase in qualified leads this quarter."
    p3b2 = tf3b.add_paragraph()
    p3b2.text = "Social media engagement rose by 18% with the new content strategy."
    p3b3 = tf3b.add_paragraph()
    p3b3.text = "Email open rates improved to 34.5%, exceeding the industry benchmark of 28%."
    for para in tf3b.paragraphs:
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.name = "Arial"

    # --- Slide 4: Q2 Projections (table) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf4t = txBox4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Q2 Projections"
    r4t = p4t.runs[0]
    r4t.font.name = "Arial"
    r4t.font.size = Pt(32)
    r4t.font.bold = True
    r4t.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    table_shape = slide4.shapes.add_table(4, 3, Inches(1), Inches(2), Inches(10), Inches(3))
    tbl = table_shape.table
    headers = ['Region', 'Revenue Target', 'Growth %']
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    rows_data = [
        ['North America', '$4.2M', '12%'],
        ['Europe', '$2.8M', '8%'],
        ['Asia Pacific', '$1.9M', '15%'],
    ]
    for r, row_data in enumerate(rows_data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # --- Slide 5: Action Items ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox5_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf5t = txBox5_title.text_frame
    p5t = tf5t.paragraphs[0]
    p5t.text = "Action Items"
    r5t = p5t.runs[0]
    r5t.font.name = "Arial"
    r5t.font.size = Pt(32)
    r5t.font.bold = True
    r5t.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)

    txBox5_body = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(5))
    tf5b = txBox5_body.text_frame
    tf5b.word_wrap = True
    items = [
        "Finalize Q3 budget allocation by June 15",
        "Launch Product B marketing campaign in APAC region",
        "Schedule customer feedback review sessions with product team",
        "Update sales training materials for new pricing structure",
        "Complete competitive analysis report for board presentation",
    ]
    tf5b.paragraphs[0].text = items[0]
    for item in items[1:]:
        p = tf5b.add_paragraph()
        p.text = item
    for para in tf5b.paragraphs:
        for run in para.runs:
            run.font.size = Pt(18)
            run.font.name = "Arial"

    prs.save(OUTPUT_PPTX)
    print(f'Presentation created: {OUTPUT_PPTX}')


def main():
    create_sales_data()
    create_presentation()

    # GUI-ready startup: open both files
    launch_gui(f'libreoffice --impress "{OUTPUT_PPTX}"', delay_sec=3.0)
    launch_gui(f'libreoffice --calc "{OUTPUT_XLSX}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress and Calc with DISPLAY=:0')


main()
