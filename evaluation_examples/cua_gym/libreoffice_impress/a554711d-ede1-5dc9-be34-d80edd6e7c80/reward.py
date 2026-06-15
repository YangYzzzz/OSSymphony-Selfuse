"""
Reward Script: Apply bold, dark blue (#00008B), and underline to titles on slides 1, 3, 5
Task ID: osworld_impress_title_selective_formatting_009
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): Titles on slides 1, 3, 5 are bold AND slides 2, 4, 6 remain not bold
  Component 2 (0.3): Titles on slides 1, 3, 5 are underlined AND slides 2, 4, 6 remain not underlined
  Component 3 (0.3): Titles on slides 1, 3, 5 have dark blue color #00008B AND slides 2, 4, 6 remain black
Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_009'

TARGET_SLIDES = [0, 2, 4]    # 0-indexed: slides 1, 3, 5 (should be formatted)
UNCHANGED_SLIDES = [1, 3, 5] # 0-indexed: slides 2, 4, 6 (should remain unmodified)
DARK_BLUE = '00008B'


def get_title_runs(slide):
    """Return all non-empty runs from the title shape of a slide."""
    runs = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith('Title'):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or '').strip():
                        runs.append(run)
    return runs


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Validate slide count (precondition gate)
    num_slides = len(prs.slides)
    if num_slides < 6:
        print(f"CRITICAL: Expected 6 slides, found {num_slides}. File may be corrupted.")
        print("REWARD: 0.0")
        return 0.0
    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: Slides 1,3,5 bold=True AND slides 2,4,6 not bold (0.4 points)
    # This compound check fails on initial (target slides not bold) and passes on golden
    try:
        target_bold_ok = True
        unchanged_bold_ok = True

        for slide_idx in TARGET_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            if not runs:
                print(f"FAIL: Component 1 — No title runs on slide {slide_idx + 1}")
                target_bold_ok = False
                continue
            if not all(run.font.bold is True for run in runs):
                actual = [run.font.bold for run in runs]
                print(f"FAIL: Component 1 — Slide {slide_idx + 1} bold={actual}, expected True")
                target_bold_ok = False
            else:
                print(f"PASS: Component 1 — Slide {slide_idx + 1} title is bold")

        for slide_idx in UNCHANGED_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            for run in runs:
                if run.font.bold is True:
                    print(f"FAIL: Component 1 — Slide {slide_idx + 1} title unexpectedly bold=True (should stay unchanged)")
                    unchanged_bold_ok = False
                    break

        if target_bold_ok and unchanged_bold_ok:
            print(f"PASS: Component 1 — Target slides (1,3,5) have bold titles; unchanged slides (2,4,6) are not bold (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 1 — Bold formatting not correctly applied/isolated")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slides 1,3,5 underline=True AND slides 2,4,6 not underlined (0.3 points)
    try:
        target_ul_ok = True
        unchanged_ul_ok = True

        for slide_idx in TARGET_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            if not runs:
                print(f"FAIL: Component 2 — No title runs on slide {slide_idx + 1}")
                target_ul_ok = False
                continue
            if not all(run.font.underline is True for run in runs):
                actual = [run.font.underline for run in runs]
                print(f"FAIL: Component 2 — Slide {slide_idx + 1} underline={actual}, expected True")
                target_ul_ok = False
            else:
                print(f"PASS: Component 2 — Slide {slide_idx + 1} title is underlined")

        for slide_idx in UNCHANGED_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            for run in runs:
                if run.font.underline is True:
                    print(f"FAIL: Component 2 — Slide {slide_idx + 1} title unexpectedly underline=True (should stay unchanged)")
                    unchanged_ul_ok = False
                    break

        if target_ul_ok and unchanged_ul_ok:
            print(f"PASS: Component 2 — Target slides (1,3,5) have underlined titles; unchanged slides (2,4,6) are not underlined (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Underline formatting not correctly applied/isolated")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 1,3,5 have color #00008B AND slides 2,4,6 still have black color (0.3 points)
    # This compound check fails on initial (target slides still black) and passes on golden
    try:
        target_color_ok = True
        unchanged_color_ok = True

        for slide_idx in TARGET_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            if not runs:
                print(f"FAIL: Component 3 — No title runs on slide {slide_idx + 1}")
                target_color_ok = False
                continue
            for run in runs:
                try:
                    if run.font.color.type is None:
                        print(f"FAIL: Component 3 — Slide {slide_idx + 1} title has no explicit color (expected {DARK_BLUE})")
                        target_color_ok = False
                        break
                    actual_color = str(run.font.color.rgb).upper()
                    if actual_color != DARK_BLUE.upper():
                        print(f"FAIL: Component 3 — Slide {slide_idx + 1} color={actual_color}, expected {DARK_BLUE}")
                        target_color_ok = False
                        break
                    else:
                        print(f"PASS: Component 3 — Slide {slide_idx + 1} title color is dark blue ({DARK_BLUE})")
                except Exception as ce:
                    print(f"FAIL: Component 3 — Slide {slide_idx + 1} color check error: {ce}")
                    target_color_ok = False
                    break

        for slide_idx in UNCHANGED_SLIDES:
            runs = get_title_runs(prs.slides[slide_idx])
            for run in runs:
                if run.font.color.type is not None:
                    try:
                        actual_color = str(run.font.color.rgb).upper()
                        if actual_color == DARK_BLUE.upper():
                            print(f"FAIL: Component 3 — Slide {slide_idx + 1} title unexpectedly has dark blue color (should remain black)")
                            unchanged_color_ok = False
                            break
                    except Exception:
                        pass  # non-RGB color types are fine for unchanged slides

        if target_color_ok and unchanged_color_ok:
            print(f"PASS: Component 3 — Target slides (1,3,5) have dark blue color; unchanged slides (2,4,6) remain black (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — Color not correctly applied/isolated")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in this env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
