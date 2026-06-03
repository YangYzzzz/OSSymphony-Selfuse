"""
Initial Setup: Quarterly Business Review Presentation (6 slides, unformatted titles)
Task ID: osworld_impress_title_selective_formatting_009
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_selective_formatting_009'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, content lines)
    slides_data = [
        (
            "Q1 2025 Business Review",
            [
                "Revenue: $4.2M (+12% YoY)",
                "New Customers: 138",
                "Customer Retention: 91%",
                "Net Promoter Score: 72",
                "Key Win: Partnership with Apex Corp",
            ]
        ),
        (
            "Market Expansion Update",
            [
                "APAC region: 3 new offices opened",
                "EMEA pipeline: $8.7M qualified",
                "Latin America: pilot program launched",
                "Total addressable market grew 18%",
                "Competitive win rate: 64%",
            ]
        ),
        (
            "Product Development Milestones",
            [
                "v3.4 released Feb 28 — 200+ features",
                "Mobile app: 4.7 stars (App Store)",
                "API integrations: 45 partners",
                "Uptime SLA: 99.97% achieved",
                "R&D investment: $1.1M this quarter",
            ]
        ),
        (
            "Human Resources & Culture",
            [
                "Headcount: 312 employees (+24 QoQ)",
                "Engineering team grew 18%",
                "eNPS (Employee NPS): 68",
                "Training hours per employee: 14",
                "Diversity hiring: 47% non-majority",
            ]
        ),
        (
            "Financial Highlights",
            [
                "Gross margin: 68.4%",
                "Operating expenses: $2.9M",
                "EBITDA: $1.3M",
                "Cash runway: 28 months",
                "ARR: $16.8M (+22% YoY)",
            ]
        ),
        (
            "Strategic Outlook Q2 2025",
            [
                "Target revenue: $4.8M",
                "Planned headcount additions: 15",
                "Major product launch: June 10",
                "Key focus: enterprise segment",
                "Board review date: April 22",
            ]
        ),
    ]

    for idx, (title_text, bullet_lines) in enumerate(slides_data):
        # Use layout 1 = Title + Content
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Set title — regular weight, black, no underline
        title_shape = slide.shapes.title
        title_shape.text = ""
        tf = title_shape.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = title_text
        run.font.bold = False
        run.font.italic = False
        run.font.underline = False
        run.font.size = Pt(36)
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)  # Black

        # Set content placeholder
        content_placeholder = slide.placeholders[1]
        tf_content = content_placeholder.text_frame
        tf_content.text = ""
        for i, line in enumerate(bullet_lines):
            if i == 0:
                p = tf_content.paragraphs[0]
            else:
                p = tf_content.add_paragraph()
            p.text = line
            p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
