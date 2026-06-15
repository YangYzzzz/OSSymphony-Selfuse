"""
Reward Script: Create a timeline on slide 4 with shapes
Task ID: impress_teach_033
Domain: libreoffice_impress
Scoring:
  Component 1: Horizontal timeline line on slide 4 (0.2 pts)
  Component 2: Four ovals labeled '1900', '1925', '1950', '1975' (0.4 pts)
  Component 3: Four vertical connector lines (0.2 pts)
  Component 4: Ovals positioned alternately above/below timeline (0.2 pts)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_033'

EXPECTED_LABELS = {'1900', '1925', '1950', '1975'}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Classify shapes on slide 4 (excluding original title/textbox shapes)
    horizontal_lines = []
    vertical_lines = []
    ovals = []
    all_lines = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's an oval
            try:
                from pptx.enum.shapes import MSO_SHAPE
                if shape.auto_shape_type == MSO_SHAPE.OVAL:
                    text = shape.text.strip() if hasattr(shape, 'text') else ''
                    ovals.append({
                        'text': text,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'center_x': shape.left + shape.width // 2,
                        'center_y': shape.top + shape.height // 2,
                    })
            except Exception:
                pass

        elif shape.shape_type == 9:  # LINE type
            w = shape.width
            h = shape.height
            # Horizontal line: width >> height
            if abs(w) > abs(h) * 3 and abs(w) > 914400:  # > 1 inch wide
                horizontal_lines.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': w,
                    'height': h,
                })
            # Vertical line: height >> width
            elif abs(h) > abs(w) * 3 and abs(h) > 100000:  # has meaningful height
                vertical_lines.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': w,
                    'height': h,
                })
            all_lines.append({
                'left': shape.left,
                'top': shape.top,
                'width': w,
                'height': h,
            })

    # Component 1: Horizontal timeline line on slide 4 (0.2 points)
    try:
        has_hline = len(horizontal_lines) >= 1
        if has_hline:
            hl = horizontal_lines[0]
            print(f"PASS: Component 1 - Horizontal timeline line found "
                  f"(width={hl['width']}, pos=({hl['left']}, {hl['top']})) (0.2 pts)")
            total_score += 0.2  # if has_hline
        else:
            print(f"FAIL: Component 1 - No horizontal line found on slide 4. "
                  f"Total lines found: {len(all_lines)}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Four ovals labeled '1900', '1925', '1950', '1975' (0.4 points)
    # Award 0.1 per correct label found
    try:
        found_labels = set()
        for oval in ovals:
            if oval['text'] in EXPECTED_LABELS:
                found_labels.add(oval['text'])

        label_count = len(found_labels)
        if label_count == 4:
            print(f"PASS: Component 2 - All 4 labeled ovals found: {sorted(found_labels)} (0.4 pts)")
            total_score += 0.4  # if label_count == 4
        elif label_count > 0:
            label_score = label_count * 0.1
            print(f"PARTIAL: Component 2 - {label_count}/4 labeled ovals found: "
                  f"{sorted(found_labels)}, missing: {sorted(EXPECTED_LABELS - found_labels)} "
                  f"({label_score} pts)")
            total_score += label_score  # elif label_count > 0
        else:
            print(f"FAIL: Component 2 - No ovals with expected labels found. "
                  f"Ovals on slide: {[o['text'] for o in ovals]}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Four vertical connector lines (0.2 points)
    # Award 0.05 per vertical line (up to 4)
    try:
        vert_count = min(len(vertical_lines), 4)
        vert_score = vert_count * 0.05
        if vert_count >= 4:
            print(f"PASS: Component 3 - {len(vertical_lines)} vertical connector lines found (0.2 pts)")
            total_score += 0.2
        elif vert_count > 0:
            print(f"PARTIAL: Component 3 - {vert_count}/4 vertical lines found ({vert_score} pts)")
            total_score += vert_score
        else:
            print(f"FAIL: Component 3 - No vertical connector lines found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Ovals positioned alternately above/below the timeline (0.2 points)
    # Check that labeled ovals alternate between above and below the horizontal line
    try:
        if len(horizontal_lines) >= 1 and len(found_labels) >= 2:
            timeline_y = horizontal_lines[0]['top']

            # Get labeled ovals sorted by x position (left to right)
            labeled_ovals = sorted(
                [o for o in ovals if o['text'] in EXPECTED_LABELS],
                key=lambda o: o['left']
            )

            above_count = 0
            below_count = 0
            positions = []
            for oval in labeled_ovals:
                # Oval center relative to timeline
                oval_center_y = oval['center_y']
                if oval_center_y < timeline_y:
                    above_count += 1
                    positions.append('above')
                else:
                    below_count += 1
                    positions.append('below')

            # Check alternation: need both above and below, and they should alternate
            if above_count >= 1 and below_count >= 1:
                # Check if consecutive ovals alternate positions
                is_alternating = all(
                    positions[j] != positions[j-1] for j in range(1, len(positions))
                )

                if is_alternating and len(positions) >= 3:
                    print(f"PASS: Component 4 - Ovals alternate above/below timeline: "
                          f"{list(zip([o['text'] for o in labeled_ovals], positions))} (0.2 pts)")
                    if is_alternating:  # guarded increment
                        total_score += 0.2
                elif above_count >= 1 and below_count >= 1:
                    # Partial: at least some are above and some below
                    print(f"PARTIAL: Component 4 - Ovals have mixed positions but not strictly alternating: "
                          f"{list(zip([o['text'] for o in labeled_ovals], positions))} (0.1 pts)")
                    if above_count >= 1:  # guarded increment
                        total_score += 0.1
                else:
                    print(f"FAIL: Component 4 - Ovals not alternating: {positions}")
            else:
                print(f"FAIL: Component 4 - Ovals not on both sides of timeline. "
                      f"above={above_count}, below={below_count}")
        else:
            print(f"FAIL: Component 4 - Cannot check alternation: "
                  f"horizontal_lines={len(horizontal_lines)}, labeled_ovals={len(found_labels)}")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
