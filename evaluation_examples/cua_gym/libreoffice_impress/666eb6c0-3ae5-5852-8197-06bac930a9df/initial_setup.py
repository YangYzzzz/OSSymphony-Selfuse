"""
Initial Setup: Create a 6-slide presentation with slide 4 titled '20th Century Milestones' but empty body.
Task ID: impress_teach_033
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_033'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "History of the 20th Century"
    slide1.placeholders[1].text = "A Journey Through Time"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "The 20th century was an era of unprecedented change."
    p2 = body2.add_paragraph()
    p2.text = "From world wars to technological revolutions, humanity reshaped the globe."
    p3 = body2.add_paragraph()
    p3.text = "This presentation explores the major milestones that defined an entire century."

    # --- Slide 3: Key Events ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Key Events"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "1914-1918: World War I reshapes European borders"
    for item in [
        "1929: The Great Depression triggers global economic turmoil",
        "1939-1945: World War II and the dawn of the atomic age",
        "1969: Apollo 11 lands on the Moon",
        "1989: Fall of the Berlin Wall signals end of Cold War",
    ]:
        p = body3.add_paragraph()
        p.text = item

    # --- Slide 4: 20th Century Milestones (EMPTY — task target) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    # Add only a title text box at the top
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "20th Century Milestones"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # --- Slide 5: Cultural Shifts ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Cultural Shifts"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "The Jazz Age brought new forms of artistic expression in the 1920s."
    for item in [
        "Post-war cinema and television transformed popular culture worldwide.",
        "The counterculture movement of the 1960s challenged traditional norms.",
        "The digital revolution in the 1990s connected billions through the internet.",
    ]:
        p = body5.add_paragraph()
        p.text = item

    # --- Slide 6: Conclusion ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Conclusion"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "The 20th century laid the foundation for our modern world."
    p = body6.add_paragraph()
    p.text = "Understanding these milestones helps us navigate the challenges of the 21st century."

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
