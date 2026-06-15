"""
Reward Script: Verify grouped bar chart on slide 5 of Performance_Review.pptx
Task ID: impress_exec_047
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Chart exists on slide 5 and is COLUMN_CLUSTERED type
  Component 2 (0.20): Chart title is 'Revenue: Actual vs Target by Region'
  Component 3 (0.25): Chart data matches specified values (3 categories, 2 series)
  Component 4 (0.15): Series 0 (Actual) fill color is #003366
  Component 5 (0.15): Series 1 (Target) fill color is #CCCCCC
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_047'


def persist_app_state(domain):
    """Save any unsaved edits in LibreOffice before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.chart import XL_CHART_TYPE
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"PRECONDITION FAIL: Expected >= 5 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed

    # Find chart shape on slide 5
    chart_shape = None
    for shape in slide5.shapes:
        if shape.has_chart:
            chart_shape = shape
            break

    # Component 1: Chart exists on slide 5 and is COLUMN_CLUSTERED (0.25 points)
    try:
        if chart_shape is not None:
            chart = chart_shape.chart
            if chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                print(f"PASS: Component 1 -- COLUMN_CLUSTERED chart found on slide 5 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 -- Chart exists but type is {chart.chart_type}, expected COLUMN_CLUSTERED")
        else:
            print(f"FAIL: Component 1 -- No chart found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Early exit if no chart
    if chart_shape is None:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    chart = chart_shape.chart

    # Component 2: Chart title is 'Revenue: Actual vs Target by Region' (0.20 points)
    try:
        if chart.has_title:
            title_text = chart.chart_title.text_frame.text.strip()
            expected_title = "Revenue: Actual vs Target by Region"
            if title_text == expected_title:
                print(f"PASS: Component 2 -- Chart title matches: '{title_text}' (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 2 -- Chart title is '{title_text}', expected '{expected_title}'")
        else:
            print(f"FAIL: Component 2 -- Chart has no title")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Chart data matches specified values (0.25 points)
    # Expected: 3 categories (North America, Europe, APAC), 2 series
    # Series 0 (Actual): [25, 18, 12]
    # Series 1 (Target): [23, 20, 10]
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        num_series = len(plot.series)

        expected_categories = ['North America', 'Europe', 'APAC']
        expected_actual = [25.0, 18.0, 12.0]
        expected_target = [23.0, 20.0, 10.0]

        data_checks_passed = 0
        data_checks_total = 3  # categories, actual values, target values
        issues = []

        # Check categories
        if categories == expected_categories:
            data_checks_passed += 1
        else:
            issues.append(f"Categories: {categories} != {expected_categories}")

        # Check series count
        if num_series != 2:
            issues.append(f"Series count: {num_series} != 2")
        else:
            # Check actual values (series 0)
            actual_vals = list(plot.series[0].values)
            if actual_vals == expected_actual:
                data_checks_passed += 1
            else:
                # Also check if stored as millions (25000000)
                actual_as_millions = [v / 1_000_000 if v > 1000 else v for v in actual_vals]
                if actual_as_millions == expected_actual:
                    data_checks_passed += 1
                else:
                    issues.append(f"Actual series values: {actual_vals} != {expected_actual}")

            # Check target values (series 1)
            target_vals = list(plot.series[1].values)
            if target_vals == expected_target:
                data_checks_passed += 1
            else:
                target_as_millions = [v / 1_000_000 if v > 1000 else v for v in target_vals]
                if target_as_millions == expected_target:
                    data_checks_passed += 1
                else:
                    issues.append(f"Target series values: {target_vals} != {expected_target}")

        if data_checks_passed == data_checks_total:
            print(f"PASS: Component 3 -- Chart data correct: 3 categories, 2 series with matching values (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 -- Data issues: {'; '.join(issues)}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Series 0 (Actual) fill color is #003366 (0.15 points)
    try:
        series0 = chart.plots[0].series[0]
        fill = series0.format.fill
        if fill.type is not None:
            color_hex = str(fill.fore_color.rgb).upper()
            if color_hex == "003366":
                print(f"PASS: Component 4 -- Actual series color is #003366 (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 4 -- Actual series color is #{color_hex}, expected #003366")
        else:
            print(f"FAIL: Component 4 -- Actual series has no solid fill")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Series 1 (Target) fill color is #CCCCCC (0.15 points)
    try:
        series1 = chart.plots[0].series[1]
        fill = series1.format.fill
        if fill.type is not None:
            color_hex = str(fill.fore_color.rgb).upper()
            if color_hex == "CCCCCC":
                print(f"PASS: Component 5 -- Target series color is #CCCCCC (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 5 -- Target series color is #{color_hex}, expected #CCCCCC")
        else:
            print(f"FAIL: Component 5 -- Target series has no solid fill")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
