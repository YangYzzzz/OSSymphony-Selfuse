"""
Initial Setup: Create 8-slide Performance Review presentation with slide 5 titled
'Regional Performance' but NO chart.
Task ID: impress_exec_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_text):
    """Add a slide with Title and Content layout (index 1)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = body_text
    p.font = p.runs[0].font if p.runs else None
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with Title Only layout (index 5)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 Performance Review"
    slide1.placeholders[1].text = "Global Operations Division\nDecember 2025"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    bullets_2 = [
        "Overall revenue grew 8.2% year-over-year to $55M across all regions",
        "North America exceeded targets by $2M driven by enterprise sales growth",
        "Europe faced headwinds from currency fluctuations but maintained profitability",
        "APAC region showed strongest growth trajectory at 15% quarter-over-quarter",
        "Customer retention rate improved to 94.3%, up from 91.7% in Q3",
    ]
    for i, txt in enumerate(bullets_2):
        if i == 0:
            tf2.paragraphs[0].text = txt
        else:
            p = tf2.add_paragraph()
            p.text = txt
            p.level = 0

    # --- Slide 3: Financial Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Financial Overview"
    tf3 = slide3.placeholders[1].text_frame
    tf3.clear()
    bullets_3 = [
        "Total Revenue: $55.0M (+8.2% YoY)",
        "Gross Margin: 68.4% (up from 65.1% in Q3)",
        "Operating Expenses: $22.8M (within budget)",
        "EBITDA: $14.9M representing 27.1% margin",
        "Free Cash Flow: $9.3M, supporting ongoing R&D investment",
    ]
    for i, txt in enumerate(bullets_3):
        if i == 0:
            tf3.paragraphs[0].text = txt
        else:
            p = tf3.add_paragraph()
            p.text = txt
            p.level = 0

    # --- Slide 4: Team Highlights ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Highlights"
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()
    bullets_4 = [
        "Engineering: Launched 3 major product features ahead of schedule",
        "Sales: Closed 12 enterprise deals worth $18.5M total contract value",
        "Marketing: Brand awareness increased 22% in target demographics",
        "Customer Success: NPS score reached all-time high of 72",
        "HR: Hired 45 new team members with 98% offer acceptance rate",
    ]
    for i, txt in enumerate(bullets_4):
        if i == 0:
            tf4.paragraphs[0].text = txt
        else:
            p = tf4.add_paragraph()
            p.text = txt
            p.level = 0

    # --- Slide 5: Regional Performance (TITLE ONLY - NO CHART) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    slide5.shapes.title.text = "Regional Performance"

    # --- Slide 6: Customer Satisfaction ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Customer Satisfaction"
    tf6 = slide6.placeholders[1].text_frame
    tf6.clear()
    bullets_6 = [
        "Overall satisfaction score: 4.6/5.0 (up from 4.3 in Q3)",
        "Support ticket resolution time reduced to 2.4 hours average",
        "Self-service portal adoption increased to 67% of all inquiries",
        "Enterprise client satisfaction: 4.8/5.0 across all regions",
        "Top feedback theme: improved onboarding experience",
    ]
    for i, txt in enumerate(bullets_6):
        if i == 0:
            tf6.paragraphs[0].text = txt
        else:
            p = tf6.add_paragraph()
            p.text = txt
            p.level = 0

    # --- Slide 7: Key Initiatives ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Key Initiatives"
    tf7 = slide7.placeholders[1].text_frame
    tf7.clear()
    bullets_7 = [
        "Platform Migration: 85% complete, on track for Q1 2026 finish",
        "AI Integration: Beta launched with 200+ pilot customers",
        "APAC Expansion: New offices in Singapore and Tokyo operational",
        "Security Certification: SOC 2 Type II audit completed successfully",
        "Partner Program: Onboarded 15 new channel partners in Q4",
    ]
    for i, txt in enumerate(bullets_7):
        if i == 0:
            tf7.paragraphs[0].text = txt
        else:
            p = tf7.add_paragraph()
            p.text = txt
            p.level = 0

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Next Steps"
    tf8 = slide8.placeholders[1].text_frame
    tf8.clear()
    bullets_8 = [
        "Finalize 2026 budget allocation by January 15th",
        "Launch Phase 2 of AI integration for enterprise tier",
        "Expand European sales team by 8 additional account executives",
        "Complete platform migration and decommission legacy systems",
        "Conduct annual employee engagement survey in February",
    ]
    for i, txt in enumerate(bullets_8):
        if i == 0:
            tf8.paragraphs[0].text = txt
        else:
            p = tf8.add_paragraph()
            p.text = txt
            p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
