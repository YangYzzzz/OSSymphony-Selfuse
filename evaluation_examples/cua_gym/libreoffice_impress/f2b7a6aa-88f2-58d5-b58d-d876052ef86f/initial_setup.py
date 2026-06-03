"""
Initial Setup: Create a 10-slide Event Deck with solid white background on slide 1
Task ID: impress_design_020
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_design_020'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    # Standard widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title Slide (WHITE background - no gradient!) ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Explicitly set white background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10), Inches(2),
                "TechVision Summit 2025", font_size=44, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(2), Inches(3.8), Inches(9), Inches(1),
                "Innovation \u00b7 Collaboration \u00b7 Future", font_size=24,
                color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(3), Inches(5.2), Inches(7), Inches(1),
                "March 15\u201317, 2025  |  San Francisco Convention Center", font_size=16,
                color=RGBColor(0x77, 0x77, 0x77), alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: Agenda ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    add_textbox(slide2, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Event Agenda", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69), alignment=PP_ALIGN.LEFT)
    agenda_items = [
        "9:00 AM \u2013 Registration & Networking Breakfast",
        "10:00 AM \u2013 Opening Keynote: The Future of AI",
        "11:30 AM \u2013 Panel: Sustainable Technology",
        "1:00 PM \u2013 Lunch & Exhibition Hall",
        "2:30 PM \u2013 Workshop Sessions (Track A/B/C)",
        "4:30 PM \u2013 Fireside Chat with Industry Leaders",
        "6:00 PM \u2013 Evening Gala & Awards Ceremony",
    ]
    for i, item in enumerate(agenda_items):
        add_textbox(slide2, Inches(1.5), Inches(1.6 + i * 0.7), Inches(9), Inches(0.6),
                    item, font_size=16, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 3: Keynote Speaker ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    add_textbox(slide3, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Keynote Speaker", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.LEFT)
    add_textbox(slide3, Inches(1.5), Inches(2), Inches(8), Inches(1),
                "Dr. Elena Rodriguez", font_size=28, bold=True,
                color=RGBColor(0xBB, 0x86, 0xFC))
    add_textbox(slide3, Inches(1.5), Inches(3), Inches(8), Inches(1),
                "Chief AI Officer, Nexus Technologies", font_size=18,
                color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(slide3, Inches(1.5), Inches(4.2), Inches(8), Inches(2),
                "\"Bridging the gap between artificial intelligence and human creativity \u2013 "
                "how the next wave of generative tools will reshape entire industries.\"",
                font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA))

    # ---- Slide 4: Panel Discussion ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    add_textbox(slide4, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Panel: Sustainable Technology", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69))
    panelists = [
        "Marcus Chen \u2013 VP of Green Computing, EcoTech",
        "Dr. Amara Okafor \u2013 Director, MIT Sustainability Lab",
        "James Park \u2013 CTO, CircularAI",
        "Sofia Mendez \u2013 Head of ESG, Global Ventures Fund",
    ]
    for i, p_text in enumerate(panelists):
        add_textbox(slide4, Inches(1.5), Inches(2 + i * 0.8), Inches(9), Inches(0.7),
                    p_text, font_size=18, color=RGBColor(0x33, 0x33, 0x33))

    # ---- Slide 5: Workshop Tracks ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_textbox(slide5, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Workshop Tracks", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69))
    tracks = [
        ("Track A: Machine Learning in Production", "Room 301 \u2013 Hands-on coding session with real-world ML pipelines"),
        ("Track B: UX Design for AI Products", "Room 405 \u2013 Interactive design sprint with Figma prototypes"),
        ("Track C: Data Privacy & Compliance", "Room 210 \u2013 Regulatory frameworks workshop with case studies"),
    ]
    for i, (title, desc) in enumerate(tracks):
        add_textbox(slide5, Inches(1.5), Inches(1.8 + i * 1.5), Inches(9), Inches(0.6),
                    title, font_size=22, bold=True, color=RGBColor(0x0D, 0x1B, 0x2A))
        add_textbox(slide5, Inches(1.5), Inches(2.4 + i * 1.5), Inches(9), Inches(0.6),
                    desc, font_size=14, color=RGBColor(0x66, 0x66, 0x66))

    # ---- Slide 6: Sponsors ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    add_textbox(slide6, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Our Sponsors", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69), alignment=PP_ALIGN.CENTER)
    sponsors = [
        "Platinum: Nexus Technologies, CloudMatrix Inc.",
        "Gold: DataForge, Quantum Leap Systems, AeroTech Solutions",
        "Silver: ByteStream Analytics, NovaSpark, TerraFlow",
    ]
    for i, s in enumerate(sponsors):
        add_textbox(slide6, Inches(2), Inches(2.5 + i * 1.0), Inches(9), Inches(0.7),
                    s, font_size=18, color=RGBColor(0x44, 0x44, 0x44), alignment=PP_ALIGN.CENTER)

    # ---- Slide 7: Venue Information ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill7 = slide7.background.fill
    fill7.solid()
    fill7.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF0)
    add_textbox(slide7, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Venue & Logistics", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69))
    add_textbox(slide7, Inches(1.5), Inches(2), Inches(8), Inches(0.6),
                "San Francisco Convention Center", font_size=24, bold=True,
                color=RGBColor(0x0D, 0x1B, 0x2A))
    venue_details = [
        "Address: 747 Howard Street, San Francisco, CA 94103",
        "Parking: Moscone Center Garage (discounted rate for attendees)",
        "Wi-Fi: Network \u2018TechVision2025\u2019 \u2013 password provided at registration",
        "Lunch: Complimentary catering in Hall B (dietary options available)",
    ]
    for i, d in enumerate(venue_details):
        add_textbox(slide7, Inches(1.5), Inches(3.2 + i * 0.7), Inches(9), Inches(0.6),
                    d, font_size=14, color=RGBColor(0x55, 0x55, 0x55))

    # ---- Slide 8: Key Statistics ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    fill8 = slide8.background.fill
    fill8.solid()
    fill8.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    add_textbox(slide8, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Event Highlights", font_size=36, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    stats = [
        ("2,500+", "Registered Attendees"),
        ("85", "Expert Speakers"),
        ("40+", "Workshop Sessions"),
        ("120", "Exhibiting Companies"),
    ]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.5 + i * 2.8)
        add_textbox(slide8, x, Inches(2.5), Inches(2.5), Inches(1),
                    num, font_size=40, bold=True,
                    color=RGBColor(0xBB, 0x86, 0xFC), alignment=PP_ALIGN.CENTER)
        add_textbox(slide8, x, Inches(3.8), Inches(2.5), Inches(0.6),
                    label, font_size=14,
                    color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    # ---- Slide 9: Networking ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    fill9 = slide9.background.fill
    fill9.solid()
    fill9.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xFF)
    add_textbox(slide9, Inches(1), Inches(0.5), Inches(10), Inches(1),
                "Networking Opportunities", font_size=36, bold=True,
                color=RGBColor(0x2D, 0x1B, 0x69))
    networking = [
        "Morning Coffee Meet & Greet (9:00 \u2013 9:45 AM)",
        "Speed Networking Lunch Tables (1:00 \u2013 1:45 PM)",
        "Industry Mixer at the Rooftop Lounge (5:30 \u2013 6:30 PM)",
        "VIP Dinner Reception (7:30 PM, invitation only)",
    ]
    for i, n in enumerate(networking):
        add_textbox(slide9, Inches(1.5), Inches(2 + i * 1.0), Inches(9), Inches(0.7),
                    n, font_size=18, color=RGBColor(0x44, 0x44, 0x44))

    # ---- Slide 10: Thank You / Contact ----
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    fill10 = slide10.background.fill
    fill10.solid()
    fill10.fore_color.rgb = RGBColor(0x2D, 0x1B, 0x69)
    add_textbox(slide10, Inches(2), Inches(2), Inches(9), Inches(1.5),
                "Thank You!", font_size=48, bold=True,
                color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, Inches(2), Inches(4), Inches(9), Inches(1),
                "info@techvisionsummit.com  |  www.techvisionsummit.com", font_size=18,
                color=RGBColor(0xCC, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
    add_textbox(slide10, Inches(2), Inches(5.2), Inches(9), Inches(0.6),
                "#TechVision2025  |  @TechVisionSummit", font_size=14,
                color=RGBColor(0xAA, 0xAA, 0xBB), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
