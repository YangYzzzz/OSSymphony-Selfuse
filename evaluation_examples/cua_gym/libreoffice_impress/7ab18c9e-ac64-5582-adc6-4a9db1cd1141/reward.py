"""
FINAL REWARD SCRIPT - SUCCESS
Task: For my upcoming client presentation, I want to add a note on slide 4 to remind myself to focus on how our services benefit them. How can I do that without affecting what's visible to the audience?
Generated: 2025-08-07 08:43:03
Status: success
Model: o4-mini
Total Steps: 3
"""

import os
from pptx import Presentation

def verify_reminder_note(file_path):
    print("Checking task completion: Add reminder note on slide 4 for benefits focus...")
    total_score = 0.0
    max_score = 1.0

    # Requirement 1: File exists (0.2 points)
    try:
        if not os.path.exists(file_path):
            print(f"✗ File not found: {file_path}")
        else:
            print(f"✓ File exists: {file_path} (0.2 points)")
            total_score += 0.2
    except Exception as e:
        print(f"✗ Error accessing file: {e}")

    # Requirement 2: At least 4 slides (0.2 points)
    prs = None
    try:
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"Found {slide_count} slides in the presentation.")
        if slide_count >= 4:
            print("✓ Slide count requirement met (0.2 points)")
            total_score += 0.2
        else:
            print("✗ Slide count less than 4 (0 points)")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")

    # Requirement 3: Slide 4 has a notes slide (0.2 points)
    notes_text = ''
    has_notes = False
    try:
        if prs and len(prs.slides) >= 4:
            slide4 = prs.slides[3]
            if hasattr(slide4, 'has_notes_slide') and slide4.has_notes_slide:
                print("✓ Slide 4 has a notes slide (0.2 points)")
                total_score += 0.2
                notes_slide = slide4.notes_slide
                # Gather all text in notes shapes
                for shape in notes_slide.shapes:
                    if hasattr(shape, 'text'):
                        txt = shape.text.strip()
                        if txt:
                            notes_text += txt + ' '
                has_notes = True
            else:
                print("✗ Slide 4 does not have a notes slide (0 points)")
        else:
            print("✗ Cannot check slide 4, not enough slides or failed to load")
    except Exception as e:
        print(f"✗ Error checking notes slide: {e}")

    # Requirement 4: Notes contain benefit-focused reminder (0.4 points)
    if has_notes:
        lower_notes = notes_text.lower()
        required_keywords = ['benefit', 'benefits']
        found_keyword = any(kw in lower_notes for kw in required_keywords)
        print(f"Extracted notes text: '{notes_text.strip()}'")
        if found_keyword:
            print("✓ Reminder note mentions benefits (0.4 points)")
            total_score += 0.4
        else:
            print("✗ Reminder note does not mention benefits (0 points)")
    else:
        print("✗ No notes text to check for benefits keyword (0 points)")

    final_score = min(total_score, max_score)
    print(f"Total score breakdown: {total_score}/{max_score}")
    print(f"Final score returned: {final_score}")
    return final_score

if __name__ == '__main__':
    file_path = '/home/user/for_my_upcoming_client_presentation_i_want_to_add_a_note_on_slide_4_to_remind_myself_to_focus_on_how.pptx'
    reward = verify_reminder_note(file_path)
    print(f"REWARD: {reward}")
