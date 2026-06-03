"""
Initial Setup: Theater show presentation with 5 slides, no transitions on slides 2 and 3.
Task ID: impress_tm_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=None, alignment=None):
    """Helper to add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "The Enchanted Stage"
    slide1.placeholders[1].text = "A Theater Production by Riverside Players"
    # Dark background for title
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1A, 0x0A, 0x2E)
    # Make title text white
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # --- Slide 2: Act I (NO transition - this is the one to get Barn Door Open) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x2C, 0x13, 0x3E)
    add_textbox(slide2, 0.5, 0.3, 9, 1.2, "Act I: The Arrival",
                font_size=36, bold=True, color=(0xFF, 0xD7, 0x00),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide2, 1, 1.8, 8, 4.5,
                "Scene opens in the grand foyer of the Belmont Theater. "
                "Crystal chandeliers cast warm golden light across marble floors. "
                "Patrons in evening attire mingle near the coat check, their "
                "voices creating a gentle hum of anticipation.\n\n"
                "MARGARET enters stage left, clutching a weathered playbill. "
                "She pauses at the ornate double doors, her gaze sweeping "
                "across the gilded moldings and velvet curtains.",
                font_size=16, color=(0xE0, 0xE0, 0xE0))

    # --- Slide 3: Act II (NO transition - this is the one to get Barn Door Close) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3E)
    add_textbox(slide3, 0.5, 0.3, 9, 1.2, "Act II: The Revelation",
                font_size=36, bold=True, color=(0xFF, 0xD7, 0x00),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide3, 1, 1.8, 8, 4.5,
                "The stage transforms into a moonlit garden behind the theater. "
                "Stone pathways wind between rose bushes, and a fountain glistens "
                "at center stage. The orchestra shifts to a haunting waltz.\n\n"
                "THOMAS emerges from the shadows, holding a sealed envelope. "
                "He crosses to the fountain and reads aloud the letter that "
                "will change everything for the Belmont company.",
                font_size=16, color=(0xE0, 0xE0, 0xE0))

    # --- Slide 4: Intermission ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)
    add_textbox(slide4, 1, 2.5, 8, 2, "~ Intermission ~",
                font_size=44, bold=True, color=(0xFF, 0xFF, 0xFF),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide4, 2, 4.5, 6, 1.5,
                "Refreshments are available in the lobby.\n"
                "The performance will resume in 15 minutes.",
                font_size=18, color=(0xFF, 0xCC, 0xCC),
                alignment=PP_ALIGN.CENTER)

    # --- Slide 5: Finale ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x2B)
    add_textbox(slide5, 0.5, 0.3, 9, 1.2, "Act III: The Final Curtain",
                font_size=36, bold=True, color=(0xFF, 0xD7, 0x00),
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide5, 1, 1.8, 8, 4.5,
                "All cast members assemble on stage for the climactic scene. "
                "The lighting shifts from cool blue to warm amber as the truth "
                "about the Belmont legacy is finally revealed.\n\n"
                "MARGARET and THOMAS stand together at center stage. "
                "The audience holds its breath as the orchestra swells "
                "to a triumphant crescendo. The curtain falls.",
                font_size=16, color=(0xE0, 0xE0, 0xE0))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
