"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’m fine-tuning my LibreOffice Impress deck and noticed that on slide 125, Table 1 is riding a bit too high. Could you drop it so its Y-position is precisely 14.0 cm? That’ll line it up neatly along the bottom of the slide.
Generated: 2025-09-10 15:49:52
Status: success
Model: azure-o3
Total Steps: 1
"""

import os
from pptx import Presentation


def verify_table_y_position(file_path: str,
                             slide_number: int = 125,
                             target_y_cm: float = 14.0) -> float:
    """Reward-script verification for table positioning on a specific slide.

    Task requirement:
        On slide 125, the first table (Table 1) must have its top-left Y
        coordinate positioned exactly at 14.0 cm.

    Scoring rubric (progressive):
        • ≤0.05 cm deviation  → 1.0  (perfect)
        • ≤0.50 cm deviation  → 0.8
        • ≤1.00 cm deviation  → 0.6
        • ≤2.00 cm deviation  → 0.4
        • ≤3.00 cm deviation  → 0.2
        •  >3.00 cm deviation → 0.0

    Only the actual accuracy of the table’s Y-position awards points; loading
    the file or finding the slide earns no points (natural conditions).
    """

    MAX_SCORE = 1.0
    score = 0.0  # default if verification fails

    # Constant for EMU→cm conversion (1 inch = 914400 EMU, 1 inch = 2.54 cm)
    EMU_PER_CM = 360000

    # ---------- 0. Preliminary file checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    if not file_path.lower().endswith(".pptx"):
        print("✗ Unsupported file type (expected .pptx)")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 1. Locate requested slide ----------
    slide_index = slide_number - 1  # convert 1-based to 0-based index
    if slide_index >= len(prs.slides):
        print(f"✗ Slide {slide_number} not present (total slides: {len(prs.slides)})")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[slide_index]

    # ---------- 2. Find first table on the slide (Table 1) ----------
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break

    if table_shape is None:
        print(f"✗ No table found on slide {slide_number}")
        print("REWARD: 0.0")
        return 0.0

    # ---------- 3. Measure Y-position and compute deviation ----------
    y_cm = table_shape.top / EMU_PER_CM
    diff_cm = abs(y_cm - target_y_cm)

    print(f"Target Y position: {target_y_cm:.2f} cm")
    print(f"Actual Y position: {y_cm:.2f} cm  (difference: {diff_cm:.4f} cm)")

    # ---------- 4. Progressive scoring ----------
    if diff_cm <= 0.05:
        score = 1.0
    elif diff_cm <= 0.5:
        score = 0.8
    elif diff_cm <= 1.0:
        score = 0.6
    elif diff_cm <= 2.0:
        score = 0.4
    elif diff_cm <= 3.0:
        score = 0.2
    else:
        score = 0.0

    # cap at MAX_SCORE
    score = min(score, MAX_SCORE)

    print(f"Computed score based on Y-position accuracy: {score}")
    print(f"REWARD: {score}")
    return score


if __name__ == "__main__":
    # Path provided in task context
    FILE_PATH = (
        "/home/user/"
        "im_fine_tuning_my_libreoffice_impress_deck_and_noticed_that_on_slide_125_table_1_is_riding_a_bit_too_golden.pptx"
    )
    verify_table_y_position(FILE_PATH)

