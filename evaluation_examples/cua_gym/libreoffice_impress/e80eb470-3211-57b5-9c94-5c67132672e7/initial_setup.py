"""
Initial Setup: Grant proposal presentation with presenter notes
Task ID: osworld_multi_apps_impress_notes_export_011
Domain: libreoffice_impress (multi-app: Impress + Writer export task)
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_011'
OUTPUT = f'{WORKDIR}/Grant_Pitch.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide data: (title, body_text, notes_text)
    # 18 slides total, 13 with non-empty notes, 5 with empty notes (slides 3, 7, 10, 14, 17)
    slides_data = [
        # Slide 1 - Title slide (has notes)
        (
            "Advancing Neuroplasticity Research:\nA Proposal for Longitudinal Study",
            "Department of Neuroscience, Westfield University\nPrincipal Investigator: Dr. Elena Vasquez\nFunding Period: 2025–2028",
            "Welcome everyone. This proposal represents three years of preliminary work.\n  Today we'll walk through the scientific rationale,   methodology,  and expected outcomes.\nWe believe this work will meaningfully advance our understanding of neuroplasticity."
        ),
        # Slide 2 - Background (has notes)
        (
            "Background: The Problem We Address",
            "• Neuroplasticity declines with age\n• Current interventions show limited efficacy\n• Longitudinal data remains scarce\n• Gap: No large-scale multi-site study exists",
            "The background section frames why this research is urgent.\n  Key point: existing studies have sample sizes under 200 — ours targets 1,500 participants.  \nEmphasize the multi-site design as our distinguishing feature."
        ),
        # Slide 3 - Prior Work (NO notes)
        (
            "Prior Work and Literature Review",
            "• Smith et al. (2019): 18-month follow-up, n=87\n• Chen & Patel (2021): Cognitive training outcomes\n• Vasquez Lab: 6 pilot studies 2020–2023\n• Meta-analysis: Bhattacharya (2022), 34 studies",
            ""
        ),
        # Slide 4 - Research Objectives (has notes)
        (
            "Research Objectives",
            "Primary Objective:\nQuantify neural plasticity changes over 36 months\n\nSecondary Objectives:\n• Identify demographic predictors of plasticity retention\n• Compare intervention modalities\n• Develop validated plasticity index",
            "These objectives align directly with the NSF program priorities for 2025.\n   Make sure to mention that objective 3 — the plasticity index — is our novel contribution.  \nThe index will be open-source and published regardless of funding outcome."
        ),
        # Slide 5 - Hypothesis (has notes)
        (
            "Central Hypothesis",
            "H1: Structured cognitive-motor training will attenuate\nplasticity decline by ≥30% over 36 months\n\nH2: Plasticity retention correlates with baseline\ncardiovascular fitness (r ≥ 0.45)",
            "Our hypotheses are conservative and grounded in pilot data.\n  H1 comes from our 2023 pilot where we saw 28% attenuation — we set 30% as the bar.  \nH2 is more exploratory but supported by three independent meta-analyses."
        ),
        # Slide 6 - Study Design (has notes)
        (
            "Study Design Overview",
            "• Randomized Controlled Trial (RCT)\n• N = 1,500 participants, ages 50–75\n• 3 arms: Cognitive training, Motor training, Combined\n• 4 assessment points: baseline, 12, 24, 36 months\n• Primary endpoint: Plasticity Index Score",
            "The RCT design is robust — reviewers last cycle asked us to strengthen it and we have.\n   We now include an active control arm to address expectation bias.  \nSite selection rationale is in Appendix B."
        ),
        # Slide 7 - Participant Criteria (NO notes)
        (
            "Participant Inclusion & Exclusion Criteria",
            "Inclusion:\n• Age 50–75, community-dwelling\n• No diagnosed neurological conditions\n• Able to perform 10-minute walk test\n\nExclusion:\n• Current psychoactive medication\n• Active malignancy\n• Severe sensory impairment",
            ""
        ),
        # Slide 8 - Intervention Protocol (has notes)
        (
            "Intervention Protocol",
            "Cognitive Training Arm:\n• 60 min/session, 3×/week, 36 months\n• Tasks: Working memory, executive function\n\nMotor Training Arm:\n• Resistance + balance, 45 min/session\n\nCombined Arm:\n• Alternating cognitive/motor sessions",
            "Protocol fidelity is ensured through a centralized training manual and monthly site audits.\n  All trainers complete a 40-hour certification — this is more rigorous than any comparable trial.  \nWe have backup trainers at each site to minimize attrition effects."
        ),
        # Slide 9 - Measurement Tools (has notes)
        (
            "Outcome Measurement Tools",
            "Neuroimaging:\n• fMRI (connectivity mapping)\n• DTI (white matter integrity)\n\nCognitive Batteries:\n• Cambridge Neuropsychological Test (CANTAB)\n• Trail Making Test A & B\n\nPhysical Assessments:\n• VO2 max, grip strength, gait analysis",
            "The neuroimaging protocol was developed with Dr. Tanaka at Stanford.\n   fMRI and DTI give us complementary views — structural and functional plasticity.  \nAll MRI sites use identical 3T Siemens scanners for comparability."
        ),
        # Slide 10 - Data Management (NO notes)
        (
            "Data Management and Quality Control",
            "• REDCap database with automated QC checks\n• Central data coordinating center (DCC) at Westfield\n• Weekly data quality reports\n• DSMB oversight every 6 months\n• Data sharing plan per NIH policy",
            ""
        ),
        # Slide 11 - Analysis Plan (has notes)
        (
            "Statistical Analysis Plan",
            "Primary Analysis:\n• Linear mixed-effects model (LME)\n• Intent-to-treat population\n• FDR correction for multiple comparisons\n\nSensitivity Analyses:\n• Per-protocol population\n• Multiple imputation for missing data",
            "Our statistician is Dr. Okonkwo — she designed the LME framework that won the ASA award.\n  The sample size was calculated assuming 20% attrition; we have dropout mitigation built in.  \nAll analysis code will be pre-registered on OSF before data collection begins."
        ),
        # Slide 12 - Power Calculation (has notes)
        (
            "Power Calculation and Sample Size",
            "• Power: 80% to detect 30% difference\n• Alpha: 0.05 (two-tailed)\n• Expected attrition: 20% over 36 months\n• Calculated N: 1,200 evaluable\n• Enrolled target: 1,500 (25% buffer)",
            "We're well-powered even under pessimistic attrition assumptions.\n   The 1,500 enrollment target gives us flexibility to run pre-specified subgroup analyses.  \nSubgroup analyses include age band, sex, and education level."
        ),
        # Slide 13 - Preliminary Results (has notes)
        (
            "Preliminary Findings from Pilot Study",
            "Pilot (n=94, 12 months):\n• Combined arm: 28% plasticity attenuation vs control\n• Effect size Cohen's d = 0.62\n• Retention rate: 91% at 12 months\n• No serious adverse events",
            "These pilot numbers are the foundation of our full proposal.\n  The 91% retention rate far exceeds typical rates of 70–75% in this age group.  \nWe attribute this to our community health worker engagement model — a key innovation."
        ),
        # Slide 14 - Sites and Partners (NO notes)
        (
            "Research Sites and Collaborating Partners",
            "Site 1: Westfield University Medical Center (Lead)\nSite 2: Pacific Coast Health System, San Diego\nSite 3: Lakewood Medical Institute, Chicago\nSite 4: Atlantic Neuroscience Center, Boston\n\nPartners: AARP, Alzheimer's Association",
            ""
        ),
        # Slide 15 - Team Qualifications (has notes)
        (
            "Research Team Qualifications",
            "PI: Dr. Elena Vasquez — 18 years neuroplasticity research\nCo-I: Dr. James Osei — clinical trial expertise\nBiostatistician: Dr. Ada Okonkwo — longitudinal methods\nProject Director: Ms. Priya Sharma — trial management\nClinical Coordinator: Mr. Rafael Diaz — community engagement",
            "Our team has collectively managed over 40 million dollars in federal grants.\n   Dr. Vasquez and Dr. Osei have co-led two prior multi-site RCTs that completed on time and on budget.  \nThe community engagement track record is especially strong — Rafael has 12 years with AARP."
        ),
        # Slide 16 - Budget Overview (has notes)
        (
            "Budget Overview (36 Months)",
            "Personnel: $2,140,000 (58%)\nEquipment & Supplies: $380,000 (10%)\nParticipant Costs: $540,000 (15%)\nTravel & Dissemination: $120,000 (3%)\nIndirect Costs: $520,000 (14%)\n\nTotal Request: $3,700,000",
            "Budget justification is in Section 4 of the proposal.\n  All personnel costs use current NIH salary caps.  \nThe participant incentive structure was designed to maximize retention — $50/visit plus transportation reimbursement."
        ),
        # Slide 17 - Timeline (NO notes)
        (
            "Project Timeline",
            "Year 1 (2025): Site setup, staff training, enrollment begins\nYear 2 (2026): Full enrollment, 12-month assessments\nYear 3 (2027): 24-month assessments, preliminary analysis\nYear 4 (2028): 36-month assessments, final analysis, dissemination\n\nMilestones tracked via Gantt chart (Appendix C)",
            ""
        ),
        # Slide 18 - Conclusion (has notes)
        (
            "Conclusion and Expected Impact",
            "• First large-scale multi-site RCT on neuroplasticity retention\n• Will generate a validated, open-source Plasticity Index\n• Directly informs clinical guidelines for aging populations\n• Training program materials freely available post-study\n• Estimated reach: 12M Americans aged 50–75",
            "In closing, this proposal addresses a critical gap with a rigorous, well-powered design.\n  The impact extends beyond our lab — the open-source Plasticity Index will serve the field for decades.  \nWe are committed to community-centered science and transparent reporting. Thank you."
        ),
    ]

    for idx, (title, body, notes) in enumerate(slides_data):
        if idx == 0:
            slide_layout = prs.slide_layouts[0]  # Title Slide
        else:
            slide_layout = prs.slide_layouts[1]  # Title + Content

        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Set body/content
        if len(slide.placeholders) > 1:
            try:
                content_ph = slide.placeholders[1]
                content_ph.text = body
            except Exception:
                pass

        # Set notes (some intentionally have extra whitespace for the agent to clean)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open Grant_Pitch.pptx in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
