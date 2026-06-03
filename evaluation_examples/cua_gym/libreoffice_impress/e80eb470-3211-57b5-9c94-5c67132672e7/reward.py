"""
Reward Script: Extract presenter notes from Grant_Pitch.pptx and save as grant_notes.docx
Task ID: osworld_multi_apps_impress_notes_export_011
Domain: libreoffice_impress (multi-app: Impress + Writer)

Scoring Rubric:
  Component 1: grant_notes.docx exists on Desktop                     — 0.3 points
  Component 2: First line reads "Slides with notes: 13"               — 0.3 points
  Component 3: All 13 non-empty notes present and whitespace-cleaned  — 0.4 points
  Total: 1.0
"""

import os
import re

WORKDIR = '/home/user/Desktop'
TASK_ID = 'osworld_multi_apps_impress_notes_export_011'

PPTX_PATH = os.path.join(WORKDIR, 'Grant_Pitch.pptx')
DOCX_PATH = os.path.join(WORKDIR, 'grant_notes.docx')

EXPECTED_SLIDES_WITH_NOTES = 13
EXPECTED_HEADER = 'Slides with notes: 13'


def clean_note_text(text):
    """
    Clean note text by stripping each line and collapsing internal spaces.
    This mirrors the expected 'clean up redundant whitespace' task requirement.
    """
    lines = text.splitlines()
    cleaned_lines = []
    for line in lines:
        # Collapse multiple internal spaces/tabs to one space, strip edges
        cleaned = ' '.join(line.split())
        cleaned_lines.append(cleaned)
    result = '\n'.join(cleaned_lines).strip()
    return result


def get_expected_notes_from_pptx():
    """
    Extract and clean all non-empty notes from the PPTX in slide order.
    Returns a list of (slide_number, cleaned_note_text) tuples.
    """
    from pptx import Presentation
    prs = Presentation(PPTX_PATH)
    notes_list = []
    for i, slide in enumerate(prs.slides):
        try:
            raw_text = slide.notes_slide.notes_text_frame.text
            cleaned = clean_note_text(raw_text)
            if cleaned:
                notes_list.append((i + 1, cleaned))
        except Exception:
            pass
    return notes_list


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Component 1: grant_notes.docx exists on Desktop (0.3 points)
    # This FAILS on initial (no docx), PASSES on golden (docx created)
    try:
        if not os.path.exists(DOCX_PATH):
            print(f"FAIL: Component 1 — grant_notes.docx not found at {DOCX_PATH}")
            print(f"\nScore: {total_score}/1.0")
            print(f"REWARD: {total_score}")
            return total_score
        print(f"PASS: Component 1 — grant_notes.docx exists at {DOCX_PATH} (0.3 pts)")
        total_score += 0.3
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Load the docx for subsequent checks
    try:
        from docx import Document
        doc = Document(DOCX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load docx {DOCX_PATH}: {e}")
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: First paragraph is "Slides with notes: 13" (0.3 points)
    # This FAILS on initial (no docx), PASSES on golden (docx has correct header)
    try:
        paragraphs = doc.paragraphs
        if len(paragraphs) == 0:
            print(f"FAIL: Component 2 — document has no paragraphs")
        else:
            first_para_text = paragraphs[0].text.strip()
            if first_para_text == EXPECTED_HEADER:
                print(f"PASS: Component 2 — first line is '{first_para_text}' (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 — expected first line '{EXPECTED_HEADER}', got '{first_para_text}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All 13 non-empty notes are present in slide order, whitespace-cleaned (0.4 points)
    # This FAILS on initial (no docx), PASSES on golden (all notes extracted and cleaned)
    try:
        # Get the expected cleaned notes from the PPTX
        expected_notes = get_expected_notes_from_pptx()

        if len(expected_notes) != EXPECTED_SLIDES_WITH_NOTES:
            print(f"WARN: Component 3 — PPTX has {len(expected_notes)} notes, expected {EXPECTED_SLIDES_WITH_NOTES}")

        # Extract non-empty paragraphs from the docx (skipping header and empty separators)
        all_para_texts = [p.text for p in doc.paragraphs]
        # Skip the first paragraph (header) and collect non-empty note paragraphs
        note_paragraphs = [t for t in all_para_texts[1:] if t.strip()]

        # Count how many expected notes are found in the docx
        found_count = 0
        missing_notes = []

        for slide_num, expected_text in expected_notes:
            # Check if this note text appears in the docx paragraphs
            note_found = False
            for para_text in note_paragraphs:
                # Compare after stripping extra whitespace from docx para too
                cleaned_para = clean_note_text(para_text)
                if cleaned_para == expected_text:
                    note_found = True
                    break
            if note_found:
                found_count += 1
            else:
                missing_notes.append(slide_num)

        # Score proportionally: partial credit for partially correct notes
        if found_count == len(expected_notes):
            print(f"PASS: Component 3 — all {found_count}/{len(expected_notes)} notes found, whitespace-cleaned (0.4 pts)")
            total_score += 0.4
        elif found_count > 0:
            partial = round(0.4 * found_count / len(expected_notes), 4)
            print(f"PARTIAL: Component 3 — {found_count}/{len(expected_notes)} notes found correctly (partial: {partial} pts)")
            print(f"  Missing from slides: {missing_notes}")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — no notes matched correctly ({found_count}/{len(expected_notes)})")
            if missing_notes:
                print(f"  Missing from slides: {missing_notes}")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {round(total_score, 4)}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


if __name__ == '__main__':
    verify_task()
