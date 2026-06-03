"""
Reward Script: Create SWOT Diagram in LibreOffice Impress from PDF Source
Task ID: pdf_cross_066
Domain: libreoffice_impress (cross-domain: pdf source)
Scoring:
  - Component 1: PPTX file exists at ~/Documents/swot_slide.pptx              (precondition gate)
  - Component 2: 4 colored SWOT quadrant shapes present (green/blue/orange/red) (0.35 pts)
  - Component 3: SWOT text content in all 4 quadrants (3-4 bullets each)       (0.40 pts)
  - Component 4: PDF export exists at ~/Documents/swot_slide.pdf               (0.15 pts)
  - Component 5: PDF export contains all SWOT section headers                  (0.10 pts)
Total: 1.0
"""

import os

WORKDIR = '/home/user'
PPTX_PATH = f'{WORKDIR}/Documents/swot_slide.pptx'
PDF_PATH = f'{WORKDIR}/Documents/swot_slide.pdf'

# Expected SWOT content keywords (case-insensitive partial match)
# Derived from task context ground truth
STRENGTHS_KEYWORDS = ['strong brand', 'loyal customer', 'skilled']
WEAKNESSES_KEYWORDS = ['online presence', 'high cost', 'aging infrastructure']
OPPORTUNITIES_KEYWORDS = ['emerging market', 'digital transform', 'partnership', 'sustainab']
THREATS_KEYWORDS = ['competitor', 'regulat', 'economic uncertainty']

# Expected colors for SWOT quadrants (with tolerance)
# green=00B050, blue=4472C4, orange=ED7D31, red=FF0000
# We accept approximate shades for each color category
def is_green(rgb_hex):
    """Check if color is a green shade."""
    r = int(rgb_hex[0:2], 16)
    g = int(rgb_hex[2:4], 16)
    b = int(rgb_hex[4:6], 16)
    # Green: g significantly > r and g significantly > b
    return g > 100 and g > r and g > b

def is_blue(rgb_hex):
    """Check if color is a blue shade."""
    r = int(rgb_hex[0:2], 16)
    g = int(rgb_hex[2:4], 16)
    b = int(rgb_hex[4:6], 16)
    # Blue: b > r or strong blue dominance
    return b > 100 and (b >= g or (b > 100 and r < 200))

def is_orange(rgb_hex):
    """Check if color is an orange shade."""
    r = int(rgb_hex[0:2], 16)
    g = int(rgb_hex[2:4], 16)
    b = int(rgb_hex[4:6], 16)
    # Orange: high r, medium-high g, low b
    return r > 150 and g > 50 and b < 100 and r > g

def is_red(rgb_hex):
    """Check if color is a red shade."""
    r = int(rgb_hex[0:2], 16)
    g = int(rgb_hex[2:4], 16)
    b = int(rgb_hex[4:6], 16)
    # Red: high r, low g and b
    return r > 150 and g < 100 and b < 100


def text_contains_keywords(text, keywords):
    """Check if text contains at least one of the keywords (case-insensitive)."""
    text_lower = text.lower()
    return any(kw.lower() in text_lower for kw in keywords)


def verify_task():
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition gate: PPTX file must exist
    if not os.path.exists(PPTX_PATH):
        print(f"FAIL (gate): PPTX file not found at {PPTX_PATH}")
        print("REWARD: 0.0")
        return 0.0

    # Load PPTX
    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(PPTX_PATH)
    except Exception as e:
        print(f"CRITICAL: Cannot load PPTX {PPTX_PATH}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("FAIL (gate): PPTX has no slides")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Collect all shapes with solid fill colors and text content
    colored_text_shapes = []
    for shape in slide.shapes:
        try:
            if shape.has_text_frame and shape.fill.type is not None:
                from pptx.enum.dml import MSO_THEME_COLOR
                fill = shape.fill
                # Check for solid fill
                if str(fill.type) == 'SOLID (1)':
                    try:
                        color_rgb = str(fill.fore_color.rgb)
                        text = shape.text_frame.text.strip()
                        if text and color_rgb and color_rgb.upper() != 'FFFFFF':
                            colored_text_shapes.append({
                                'color': color_rgb,
                                'text': text,
                                'name': shape.name
                            })
                    except Exception:
                        pass
        except Exception:
            pass

    # Component 2: 4 SWOT quadrant shapes with correct colors (0.35 pts)
    # Need 4 colored non-white shapes each with a different SWOT color category
    try:
        color_categories_found = {'green': False, 'blue': False, 'orange': False, 'red': False}
        for s in colored_text_shapes:
            c = s['color'].upper()
            if len(c) == 6:
                if is_green(c):
                    color_categories_found['green'] = True
                if is_blue(c):
                    color_categories_found['blue'] = True
                if is_orange(c):
                    color_categories_found['orange'] = True
                if is_red(c):
                    color_categories_found['red'] = True

        colors_found_count = sum(color_categories_found.values())
        if colors_found_count == 4:
            print(f"PASS: Component 2 — All 4 SWOT quadrant colors found (green/blue/orange/red) (0.35 pts)")
            total_score += 0.35
        elif colors_found_count >= 2:
            partial = 0.35 * (colors_found_count / 4)
            print(f"PARTIAL: Component 2 — {colors_found_count}/4 SWOT colors found: {color_categories_found} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 — Expected 4 SWOT colors (green/blue/orange/red), found {colors_found_count}: {color_categories_found}")
            print(f"  Colored shapes detected: {[(s['name'], s['color']) for s in colored_text_shapes]}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: SWOT text content in all 4 quadrants (0.40 pts)
    # Check that each quadrant section contains the expected keywords
    try:
        all_text = ' '.join(s['text'] for s in colored_text_shapes)
        # Also include all shapes text (in case color extraction missed some)
        full_text = ' '.join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)

        sections_found = {
            'strengths': False,
            'weaknesses': False,
            'opportunities': False,
            'threats': False
        }

        # Check each section header and its keywords
        full_text_lower = full_text.lower()
        if 'strength' in full_text_lower and text_contains_keywords(full_text, STRENGTHS_KEYWORDS):
            sections_found['strengths'] = True
        if 'weakness' in full_text_lower and text_contains_keywords(full_text, WEAKNESSES_KEYWORDS):
            sections_found['weaknesses'] = True
        if 'opportunit' in full_text_lower and text_contains_keywords(full_text, OPPORTUNITIES_KEYWORDS):
            sections_found['opportunities'] = True
        if 'threat' in full_text_lower and text_contains_keywords(full_text, THREATS_KEYWORDS):
            sections_found['threats'] = True

        sections_count = sum(sections_found.values())
        if sections_count == 4:
            print(f"PASS: Component 3 — All 4 SWOT sections with correct content found (0.40 pts)")
            total_score += 0.40
        elif sections_count >= 1:
            partial = 0.40 * (sections_count / 4)
            print(f"PARTIAL: Component 3 — {sections_count}/4 SWOT sections with content: {sections_found} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No SWOT section content found in slide shapes")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: PDF export exists at ~/Documents/swot_slide.pdf (0.15 pts)
    try:
        if os.path.exists(PDF_PATH) and os.path.getsize(PDF_PATH) > 0:
            print(f"PASS: Component 4 — PDF export exists at {PDF_PATH} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 — PDF export not found at {PDF_PATH}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: PDF export contains all SWOT section headers (0.10 pts)
    try:
        if os.path.exists(PDF_PATH):
            try:
                import pymupdf
            except ImportError:
                import fitz as pymupdf

            doc = pymupdf.open(PDF_PATH)
            pdf_text = ''
            for page in doc:
                pdf_text += page.get_text()
            doc.close()

            pdf_text_lower = pdf_text.lower()
            swot_headers_in_pdf = (
                'strength' in pdf_text_lower and
                'weakness' in pdf_text_lower and
                'opportunit' in pdf_text_lower and
                'threat' in pdf_text_lower
            )
            if swot_headers_in_pdf:
                print(f"PASS: Component 5 — PDF export contains all SWOT headers (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 — PDF export does not contain all 4 SWOT headers")
                print(f"  PDF text preview: {repr(pdf_text[:300])}")
        else:
            print(f"FAIL: Component 5 — PDF file not found, cannot verify content")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


verify_task()
