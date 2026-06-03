"""
FINAL REWARD SCRIPT - SUCCESS
Task: Make the first sentence of paragraph 4 bold and blue.
Generated: 2025-10-17 08:23:55
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

BLUE = RGBColor(0, 0, 255)  # Exact blue required (RGB 0,0,255)

# ---------------- Helper Functions -----------------

def safe_rgb(run):
    """Return RGBColor for a run if explicitly set to RGB color; otherwise None."""
    col = run.font.color
    if col is None:
        return None
    try:
        return col.rgb  # Accessible only when color.type == RGB
    except Exception:
        return None

def paragraph_text(paragraph):
    """Concatenate all run texts in a paragraph."""
    return "".join(run.text for run in paragraph.runs)

def extract_non_empty_paragraphs(presentation):
    """Return list of non-empty paragraphs (title + body) in order of appearance."""
    paragraphs = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                if paragraph_text(para).strip():
                    paragraphs.append(para)
    return paragraphs

# ---------------- Verification Logic ---------------

def verify_first_sentence_bold_blue(file_path):
    """Progressively score whether the first sentence of paragraph 4 is bold AND blue."""
    print(f"Verifying file: {file_path}")
    max_score = 1.0
    score = 0.0

    # 1. Load presentation (prerequisite – no points)
    if not os.path.exists(file_path):
        print("✗ File not found")
        return 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # 2. Collect non-empty paragraphs
    paragraphs = extract_non_empty_paragraphs(prs)
    print(f"Total non-empty paragraphs (incl. title): {len(paragraphs)}")

    # Expect: title paragraph + 4 body paragraphs = 5 non-empty paragraphs
    if len(paragraphs) >= 5:
        body_paragraphs = paragraphs[1:]  # skip title
    else:
        print("✗ Fewer than 5 non-empty paragraphs – cannot locate 4th paragraph")
        return 0.0

    # 3. Verify at least 4 body paragraphs exist (0.4 pts)
    if len(body_paragraphs) >= 4:
        score += 0.4
        print("✓ Found ≥4 body paragraphs (0.4 points)")
    else:
        print("✗ Less than 4 body paragraphs")
        return score

    # 4. Target: 4th body paragraph (index 3) – get its first sentence runs
    target_para = body_paragraphs[3]
    full_text = paragraph_text(target_para).strip()
    print(f"Paragraph 4 text: '{full_text}'")

    first_sentence_runs = []
    sentence_end_found = False
    for run in target_para.runs:
        first_sentence_runs.append(run)
        if any(ch in run.text for ch in ".!?"):
            sentence_end_found = True
            break
    if not sentence_end_found:
        print("! No terminating punctuation – entire paragraph treated as first sentence")

    meaningful_runs = [r for r in first_sentence_runs if r.text.strip()]
    if not meaningful_runs:
        print("✗ No text found in first sentence")
        return score

    # 5. Bold verification (0.3 pts)
    bold_ok = all(getattr(r.font, "bold", False) is True for r in meaningful_runs)
    if bold_ok:
        score += 0.3
        print("✓ First-sentence text is bold (0.3 points)")
    else:
        print("✗ First-sentence text is not fully bold")

    # 6. Blue color verification (0.3 pts)
    blue_ok = all(safe_rgb(r) == BLUE for r in meaningful_runs)
    if blue_ok:
        score += 0.3
        print("✓ First-sentence text is blue (0.3 points)")
    else:
        print("✗ First-sentence text is not fully blue")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------- Entrypoint -----------------------
if __name__ == "__main__":
    verify_first_sentence_bold_blue("/home/user/make_the_first_sentence_of_paragraph_4_bold_and_blue.pptx")

