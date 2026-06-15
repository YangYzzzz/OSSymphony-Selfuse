"""
Reward Script: Hyperlink navigation on agenda slide and back links
Task ID: impress_rp_038
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): 5 agenda items on slide 2 each hyperlink to correct target slide
  Component 2 (0.3): Slides 3,5,7,9,11 each have a "Back to Agenda" text shape
  Component 3 (0.2): Each "Back to Agenda" shape hyperlinks back to slide 2
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_038'

# Mapping: agenda item index (0-based among 5 items) -> target slide number
AGENDA_TARGETS = {
    0: 3,   # Item 1 -> slide 3
    1: 5,   # Item 2 -> slide 5
    2: 7,   # Item 3 -> slide 7
    3: 9,   # Item 4 -> slide 9
    4: 11,  # Item 5 -> slide 11
}

BACK_LINK_SLIDES = [3, 5, 7, 9, 11]  # slides that should have "Back to Agenda"
AGENDA_SLIDE_NUM = 2  # slide number for agenda


def get_slide_hyperlink_targets(pptx_path, slide_num):
    """
    Parse the relationships file for a given slide to find all slide-to-slide hyperlinks.
    Returns dict: {rId: target_slide_filename} e.g. {'rId2': 'slide3.xml'}
    """
    rels = {}
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            rels_path = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
            with zf.open(rels_path) as f:
                root = ET.parse(f).getroot()
                for rel in root:
                    target = rel.get('Target', '')
                    rid = rel.get('Id', '')
                    rtype = rel.get('Type', '')
                    # Only internal slide links
                    if 'slide' in rtype.split('/')[-1].lower() and target.startswith('slide'):
                        rels[rid] = target
    except Exception as e:
        print(f"WARN: Could not parse rels for slide {slide_num}: {e}")
    return rels


def get_hlinkclick_elements(pptx_path, slide_num):
    """
    Parse slide XML to find all hlinkClick elements with their associated text.
    Returns list of (text, rId, action) tuples.
    """
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    }
    results = []
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            with zf.open(f'ppt/slides/slide{slide_num}.xml') as f:
                root = ET.parse(f).getroot()
                # Find all shape trees
                for sp in root.findall('.//p:sp', ns):
                    # Collect all text and hyperlinks from this shape
                    shape_text = ""
                    shape_links = []
                    for txBody in sp.findall('.//p:txBody', ns):
                        for para in txBody.findall('a:p', ns):
                            for run in para.findall('a:r', ns):
                                t = run.find('a:t', ns)
                                run_text = t.text if t is not None else ""
                                shape_text += run_text if run_text else ""
                                rPr = run.find('a:rPr', ns)
                                if rPr is not None:
                                    hl = rPr.find('a:hlinkClick', ns)
                                    if hl is not None:
                                        rid = hl.get(f'{{{ns["r"]}}}id', '')
                                        action = hl.get('action', '')
                                        shape_links.append((run_text, rid, action))
                    for link_text, rid, action in shape_links:
                        results.append((shape_text.strip(), link_text, rid, action))
    except Exception as e:
        print(f"WARN: Could not parse slide {slide_num} XML: {e}")
    return results


def verify_task(file_path):
    """
    Verify hyperlink navigation setup with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file is a valid pptx
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        if len(prs.slides) < 11:
            print(f"CRITICAL: Expected at least 11 slides, found {len(prs.slides)}")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot load file: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Agenda items on slide 2 hyperlink to correct target slides (0.5 points)
    # Each of the 5 links is worth 0.1 points
    try:
        slide2_rels = get_slide_hyperlink_targets(file_path, 2)
        slide2_links = get_hlinkclick_elements(file_path, 2)

        # We expect 5 hyperlinks with action ppaction://hlinksldjump pointing to correct slides
        agenda_link_count = 0
        for shape_text, run_text, rid, action in slide2_links:
            if 'hlinksldjump' in action and rid in slide2_rels:
                target_file = slide2_rels[rid]
                # Extract target slide number from filename like "slide3.xml"
                try:
                    target_num = int(target_file.replace('slide', '').replace('.xml', ''))
                except ValueError:
                    continue

                # Check if this matches one of our expected targets
                expected_targets = list(AGENDA_TARGETS.values())
                if target_num in expected_targets:
                    agenda_link_count += 1
                    print(f"PASS: Agenda link found -> slide {target_num} (text: '{run_text[:40]}')")
                else:
                    print(f"INFO: Agenda link to unexpected slide {target_num}")

        # Deduplicate: count unique target slides linked
        linked_targets = set()
        for shape_text, run_text, rid, action in slide2_links:
            if 'hlinksldjump' in action and rid in slide2_rels:
                target_file = slide2_rels[rid]
                try:
                    target_num = int(target_file.replace('slide', '').replace('.xml', ''))
                except ValueError:
                    continue
                if target_num in AGENDA_TARGETS.values():
                    linked_targets.add(target_num)

        comp1_score = len(linked_targets) * 0.1
        total_score += comp1_score
        print(f"Component 1: {len(linked_targets)}/5 agenda hyperlinks correct ({comp1_score:.1f}/0.5 pts)")
        if len(linked_targets) < 5:
            missing = set(AGENDA_TARGETS.values()) - linked_targets
            print(f"  Missing links to slides: {sorted(missing)}")

    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: "Back to Agenda" text exists on target slides (0.3 points)
    # Each of the 5 slides is worth 0.06 points
    try:
        back_text_count = 0
        for slide_num in BACK_LINK_SLIDES:
            slide = prs.slides[slide_num - 1]  # 0-indexed
            found_back_text = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text = shape.text_frame.text.strip().lower()
                    if 'back to agenda' in full_text or ('back' in full_text and 'agenda' in full_text):
                        found_back_text = True
                        break
                # Also check grouped shapes
                if hasattr(shape, 'shapes'):
                    for sub in shape.shapes:
                        if hasattr(sub, 'text_frame') and sub.text_frame:
                            st = sub.text_frame.text.strip().lower()
                            if 'back to agenda' in st or ('back' in st and 'agenda' in st):
                                found_back_text = True
                                break

            if found_back_text:
                back_text_count += 1
                print(f"PASS: Slide {slide_num} has 'Back to Agenda' text")
            else:
                print(f"FAIL: Slide {slide_num} missing 'Back to Agenda' text")

        comp2_score = back_text_count * 0.06
        total_score += comp2_score
        print(f"Component 2: {back_text_count}/5 'Back to Agenda' texts found ({comp2_score:.2f}/0.3 pts)")

    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: "Back to Agenda" shapes hyperlink to slide 2 (0.2 points)
    # Each of the 5 slides is worth 0.04 points
    try:
        back_link_count = 0
        for slide_num in BACK_LINK_SLIDES:
            slide_rels = get_slide_hyperlink_targets(file_path, slide_num)
            slide_links = get_hlinkclick_elements(file_path, slide_num)

            found_back_link = False
            for shape_text, run_text, rid, action in slide_links:
                text_lower = shape_text.lower()
                if ('back' in text_lower and 'agenda' in text_lower) or 'back to agenda' in text_lower:
                    if 'hlinksldjump' in action and rid in slide_rels:
                        target_file = slide_rels[rid]
                        try:
                            target_num = int(target_file.replace('slide', '').replace('.xml', ''))
                        except ValueError:
                            continue
                        if target_num == AGENDA_SLIDE_NUM:
                            found_back_link = True
                            break

            if found_back_link:
                back_link_count += 1
                print(f"PASS: Slide {slide_num} 'Back to Agenda' links to slide 2")
            else:
                print(f"FAIL: Slide {slide_num} 'Back to Agenda' does not link to slide 2")

        comp3_score = back_link_count * 0.04
        total_score += comp3_score
        print(f"Component 3: {back_link_count}/5 'Back to Agenda' hyperlinks correct ({comp3_score:.2f}/0.2 pts)")

    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
