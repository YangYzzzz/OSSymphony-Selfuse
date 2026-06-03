"""
Initial Setup: Workshop Guide presentation with 15 slides, agenda on slide 2
Task ID: impress_rp_038
Domain: libreoffice_impress

Creates a 15-slide Workshop Guide presentation. Slide 2 is the agenda with
5 items. No hyperlinks exist anywhere in the presentation.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_038'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, body_lines, layout_idx=1):
    """Add a slide with title and bulleted body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_text_slide(prs, title_text, content_paragraphs):
    """Add a slide with a title and a text box with multiple paragraphs."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Title text box
    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf_title = tb_title.text_frame
    p = tf_title.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    p.alignment = PP_ALIGN.LEFT

    # Content text box
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para_text in enumerate(content_paragraphs):
        if i == 0:
            tf.paragraphs[0].text = para_text
        else:
            p2 = tf.add_paragraph()
            p2.text = para_text
        p_ref = tf.paragraphs[i]
        p_ref.font.size = Pt(16)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== Slide 1: Title Slide =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Workshop Guide: Advanced Data Analytics"
    slide1.placeholders[1].text = "Q2 2025 Training Program\nFacilitated by Dr. Elena Vasquez"

    # ===== Slide 2: Agenda =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    # Agenda title
    tb_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf_title = tb_title.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Workshop Agenda"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    p_title.alignment = PP_ALIGN.CENTER

    # 5 agenda items as separate text boxes for easy hyperlink targeting
    agenda_items = [
        "1. Data Collection & Preparation",
        "2. Exploratory Data Analysis",
        "3. Statistical Modeling Techniques",
        "4. Machine Learning Fundamentals",
        "5. Visualization & Reporting",
    ]
    y_start = Inches(1.8)
    for idx, item_text in enumerate(agenda_items):
        tb = slide2.shapes.add_textbox(
            Inches(1.5), y_start + Inches(idx * 0.9), Inches(7), Inches(0.7)
        )
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = item_text
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x2C, 0x5F, 0x8A)
        p.alignment = PP_ALIGN.LEFT

    # ===== Slide 3: Section 1 - Data Collection (target for Item 1) =====
    add_blank_text_slide(prs, "Data Collection & Preparation", [
        "Effective data collection is the cornerstone of any analytics project.",
        "Key sources include transactional databases, API feeds, IoT sensors, and survey instruments.",
        "Data preparation involves cleaning, deduplication, type casting, and handling missing values.",
        "Tools covered: Python pandas, Apache Spark, OpenRefine.",
        "Exercise: Clean the provided customer_transactions.csv dataset with 15,000 records.",
    ])

    # ===== Slide 4: Section 1 continued =====
    add_blank_text_slide(prs, "Data Quality Assessment", [
        "Metrics for data quality: completeness, consistency, accuracy, timeliness.",
        "Automated profiling with Great Expectations and Deequ libraries.",
        "Case study: How Meridian Corp reduced data errors by 42% in Q3 2024.",
        "Hands-on: Build a data quality dashboard using Grafana.",
    ])

    # ===== Slide 5: Section 2 - EDA (target for Item 2) =====
    add_blank_text_slide(prs, "Exploratory Data Analysis", [
        "EDA helps uncover patterns, anomalies, and relationships before formal modeling.",
        "Univariate analysis: distributions, outliers, central tendency.",
        "Bivariate analysis: correlations, scatter plots, contingency tables.",
        "Multivariate techniques: PCA, clustering for dimensionality reduction.",
        "Lab: Perform EDA on the retail_sales_2024.parquet dataset.",
    ])

    # ===== Slide 6: Section 2 continued =====
    add_blank_text_slide(prs, "Advanced EDA Techniques", [
        "Time series decomposition: trend, seasonality, residuals.",
        "Geographic analysis using GeoPandas and Folium.",
        "Text feature extraction with TF-IDF and word embeddings.",
        "Interactive exploration with Plotly and Streamlit dashboards.",
    ])

    # ===== Slide 7: Section 3 - Statistical Modeling (target for Item 3) =====
    add_blank_text_slide(prs, "Statistical Modeling Techniques", [
        "Hypothesis testing: t-tests, chi-squared, ANOVA for group comparisons.",
        "Regression analysis: linear, logistic, and polynomial models.",
        "Bayesian inference and its practical applications in A/B testing.",
        "Model selection criteria: AIC, BIC, cross-validation scores.",
        "Workshop: Build a pricing model for the housing_market dataset.",
    ])

    # ===== Slide 8: Section 3 continued =====
    add_blank_text_slide(prs, "Regression Deep Dive", [
        "Assumptions of linear regression: linearity, homoscedasticity, normality.",
        "Diagnosing multicollinearity with VIF scores.",
        "Regularization: Ridge, Lasso, and Elastic Net approaches.",
        "Real-world example: Predicting employee attrition at TechNova Inc.",
    ])

    # ===== Slide 9: Section 4 - ML Fundamentals (target for Item 4) =====
    add_blank_text_slide(prs, "Machine Learning Fundamentals", [
        "Supervised learning: classification and regression problem framing.",
        "Key algorithms: Decision Trees, Random Forests, Gradient Boosting, SVM.",
        "Unsupervised learning: K-Means, DBSCAN, hierarchical clustering.",
        "Evaluation metrics: accuracy, precision, recall, F1, ROC-AUC.",
        "Project: Build a customer churn prediction pipeline.",
    ])

    # ===== Slide 10: Section 4 continued =====
    add_blank_text_slide(prs, "Model Training Best Practices", [
        "Feature engineering strategies for tabular data.",
        "Train/validation/test splits and stratification techniques.",
        "Hyperparameter tuning with GridSearch, RandomSearch, and Optuna.",
        "Model interpretability: SHAP values and LIME explanations.",
    ])

    # ===== Slide 11: Section 5 - Visualization (target for Item 5) =====
    add_blank_text_slide(prs, "Visualization & Reporting", [
        "Principles of effective data visualization: clarity, accuracy, efficiency.",
        "Chart selection guide: when to use bar, line, scatter, heatmap, treemap.",
        "Dashboard design: layout hierarchy, color theory, interactivity patterns.",
        "Tools ecosystem: Matplotlib, Seaborn, Plotly, Tableau, Power BI.",
        "Capstone: Create an executive dashboard from the quarterly_report dataset.",
    ])

    # ===== Slide 12: Section 5 continued =====
    add_blank_text_slide(prs, "Storytelling with Data", [
        "Narrative structure: situation, complication, resolution framework.",
        "Annotation techniques to guide audience attention.",
        "Designing for different audiences: technical vs. executive stakeholders.",
        "Case study: How Aurora Analytics presented findings to the board.",
    ])

    # ===== Slide 13: Workshop Resources =====
    add_blank_text_slide(prs, "Workshop Resources", [
        "GitHub repository: github.com/analytics-workshop/q2-2025-materials",
        "Dataset downloads available on the internal SharePoint portal.",
        "Recommended reading: 'Practical Statistics for Data Scientists' by Bruce & Bruce.",
        "Online practice environment: JupyterHub at hub.training.meridian.internal.",
        "Office hours: Tuesdays and Thursdays, 2:00 PM - 4:00 PM EST.",
    ])

    # ===== Slide 14: Schedule & Logistics =====
    add_blank_text_slide(prs, "Schedule & Logistics", [
        "Week 1 (Apr 7-11): Data Collection & Preparation",
        "Week 2 (Apr 14-18): Exploratory Data Analysis",
        "Week 3 (Apr 21-25): Statistical Modeling Techniques",
        "Week 4 (Apr 28 - May 2): Machine Learning Fundamentals",
        "Week 5 (May 5-9): Visualization & Reporting",
        "Final presentations: May 12, 2025 at 10:00 AM EST.",
    ])

    # ===== Slide 15: Q&A / Closing =====
    slide15 = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide15.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\nContact: elena.vasquez@meridian.com"
    p2.font.size = Pt(18)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Slack: #analytics-workshop-q2"
    p3.font.size = Pt(18)
    p3.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
