"""
Initial Setup: Apply strikethrough formatting to the first bullet on slide 3
Task ID: osworld_impress_strikethrough_text_001
Domain: libreoffice_impress

Creates a 5-slide project planning presentation. Slide 3 has 4 bullet points
representing project tasks — none have strikethrough formatting in the initial state.
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_001'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Aurora: Q2 2025 Execution Plan"
    slide1.placeholders[1].text = "Prepared by the Strategy & Operations Team\nApril 2025"

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Project Overview & Goals"
    for item in [
        "Team Structure & Roles",
        "Milestones & Timeline",
        "Task Assignments",
        "Risks & Mitigations",
    ]:
        p = tf2.add_paragraph()
        p.text = item
        p.level = 0

    # ── Slide 3: Task Assignments (content slide with 4 bullets) ─────────────
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Task Assignments — Phase 1"
    tf3 = slide3.placeholders[1].text_frame
    # Four bullet items representing project tasks — NO strikethrough in initial state
    bullets = [
        "Complete stakeholder requirements analysis by April 12",
        "Finalize technical architecture and infrastructure plan",
        "Develop initial prototype for user acceptance testing",
        "Coordinate cross-team review sessions with QA and Design",
    ]
    tf3.text = bullets[0]
    for item in bullets[1:]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    # ── Slide 4: Milestones & Timeline ────────────────────────────────────────
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Key Milestones — Q2 2025"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "April 2025: Requirements lock and team onboarding"
    for item in [
        "May 2025: Prototype development and internal review",
        "June 2025: UAT launch and stakeholder sign-off",
        "End of Q2: Full deployment and post-launch monitoring",
    ]:
        p = tf4.add_paragraph()
        p.text = item
        p.level = 0

    # ── Slide 5: Risks & Mitigations ─────────────────────────────────────────
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Risks & Mitigation Strategies"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Risk: Resource availability constraints in May"
    for item in [
        "Mitigation: Cross-train team members for key roles",
        "Risk: Third-party API integration delays",
        "Mitigation: Establish fallback manual workflows",
    ]:
        p = tf5.add_paragraph()
        p.text = item
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open initial artifact in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
