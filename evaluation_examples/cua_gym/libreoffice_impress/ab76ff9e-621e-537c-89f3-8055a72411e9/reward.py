"""
Reward Script: Apply strikethrough formatting to the first line of the bullet list on slide 3.
Task ID: osworld_impress_strikethrough_text_001
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.5): First bullet on slide 3 has sngStrike formatting
  - Component 2 (0.3): Compound check — first bullet IS struck AND bullets 2-4 are NOT struck
  - Component 3 (0.2): Compound check — first bullet IS struck AND no other slide has strikes
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_strikethrough_text_001'


def persist_app_state():
    """Best-effort save hook for LibreOffice Impress in case of unsaved GUI edits."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def get_run_strike(run):
    """Return strike attribute value for a run, defaulting to 'noStrike'."""
    return run.font._element.attrib.get('strike', 'noStrike')


def count_struck_runs_in_slide(prs, slide_idx):
    """Count the number of runs with any strikethrough in a given slide (0-indexed)."""
    count = 0
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip() and get_run_strike(run) != 'noStrike':
                        count += 1
    return count


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Apply strikethrough to the first bullet line on slide 3 ONLY.

    Slide 3 (index 2), Content Placeholder 2, has 4 bullet paragraphs:
      Para 0: 'Complete stakeholder requirements analysis by April 12'   <- must be sngStrike
      Para 1: 'Finalize technical architecture and infrastructure plan'  <- must NOT be struck
      Para 2: 'Develop initial prototype for user acceptance testing'    <- must NOT be struck
      Para 3: 'Coordinate cross-team review sessions with QA and Design' <- must NOT be struck

    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition gate: verify presentation has at least 3 slides
    if len(prs.slides) < 3:
        print(f"CRITICAL: Expected at least 3 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide3 = prs.slides[2]

    # Find the content placeholder on slide 3 (the bullet list shape)
    content_shape = None
    for shape in slide3.shapes:
        if shape.has_text_frame and shape.name == 'Content Placeholder 2':
            content_shape = shape
            break

    if content_shape is None:
        print("CRITICAL: 'Content Placeholder 2' not found on slide 3")
        print("REWARD: 0.0")
        return 0.0

    paragraphs = content_shape.text_frame.paragraphs
    if len(paragraphs) < 1:
        print("CRITICAL: No paragraphs found in the content placeholder on slide 3")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: First bullet (para 0) has sngStrike on all non-empty runs (0.5 points)
    # FAILS on initial (noStrike) → PASSES on golden (sngStrike)
    try:
        first_para = paragraphs[0]
        nonempty_runs_p0 = [r for r in first_para.runs if (r.text or "").strip()]
        if not nonempty_runs_p0:
            print("FAIL: Component 1 — first bullet has no non-empty runs")
        else:
            strike_values = [get_run_strike(r) for r in nonempty_runs_p0]
            if all(s == 'sngStrike' for s in strike_values):
                print(f"PASS: Component 1 — first bullet has sngStrike on all runs (0.5 pts) "
                      f"[text: {repr(nonempty_runs_p0[0].text[:40])}]")
                total_score += 0.5
            else:
                print(f"FAIL: Component 1 — first bullet strike values: {strike_values}, "
                      f"expected all 'sngStrike'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Compound — first bullet IS struck AND bullets 2-4 are NOT struck (0.3 points)
    # Gated on Component 1: if first line not struck, skip (would be trivially true otherwise).
    # FAILS on initial (gate: first line not struck) → PASSES on golden
    try:
        if total_score >= 0.5:
            # Collect all struck runs in paragraphs 1-3
            other_struck = []
            for para_idx in range(1, min(4, len(paragraphs))):
                para = paragraphs[para_idx]
                for run in para.runs:
                    if (run.text or "").strip() and get_run_strike(run) != 'noStrike':
                        other_struck.append((para_idx, run.text[:40], get_run_strike(run)))
            if len(other_struck) == 0:
                print("PASS: Component 2 — bullets 2-4 on slide 3 have noStrike (0.3 pts)")
                total_score += 0.3
            else:
                for para_idx, text, strike in other_struck:
                    print(f"FAIL: Component 2 — para {para_idx} has strike={strike} "
                          f"[text: {repr(text)}]")
        else:
            print("SKIP: Component 2 — skipped because Component 1 did not pass "
                  "(first bullet not struck; cannot verify scope)")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Compound — first bullet IS struck AND no other slide has any strikethrough (0.2 points)
    # Gated on Component 1: if first line not struck, skip (would be trivially true otherwise).
    # FAILS on initial (gate: first line not struck) → PASSES on golden
    try:
        if total_score >= 0.5:
            # Sum struck runs across all slides except slide 3
            unexpected_strikes = []
            for slide_idx in range(len(prs.slides)):
                if slide_idx == 2:
                    continue  # skip slide 3, we already handled it
                n = count_struck_runs_in_slide(prs, slide_idx)
                if n > 0:
                    unexpected_strikes.append((slide_idx + 1, n))

            if len(unexpected_strikes) == 0:
                print("PASS: Component 3 — no strikethrough on other slides (0.2 pts)")
                total_score += 0.2
            else:
                for slide_num, n in unexpected_strikes:
                    print(f"FAIL: Component 3 — slide {slide_num} has {n} unexpected struck run(s)")
        else:
            print("SKIP: Component 3 — skipped because Component 1 did not pass")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
