"""
Reward Script: Set unique text alignment per slide
Task ID: osworld_impress_per_slide_alignment_008
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 2 all paragraphs are RIGHT-aligned (0.35 pts)
  - Component 2: Slide 3 all paragraphs are CENTER-aligned (0.35 pts)
  - Component 3: Slide 5 all paragraphs are RIGHT-aligned (0.30 pts)
  NOTE: Slide 4 (LEFT) is already left-aligned in initial_env, so it is a
        precondition and NOT a scoring component (would pass on both envs).
        Slides 1, 6, 7 are also unchanged (not scored).
"""

import os
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_008'


def get_dominant_alignment(slide):
    """
    Returns the majority alignment found across all non-empty paragraphs
    in the slide. Uses PP_ALIGN enum value integers for comparison.
    Returns None if no text found.
    """
    alignments = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    alignments.append(para.alignment)
    return alignments


def all_paragraphs_have_alignment(slide, expected_align):
    """
    Check that ALL non-empty paragraphs in the slide have the expected alignment.
    Returns (bool, int, int) — (all_match, matching_count, total_count)
    """
    total = 0
    matching = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    total += 1
                    if para.alignment == expected_align:
                        matching += 1
    return (matching == total and total > 0), matching, total


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Set unique text alignment per slide:
      - Slide 2 = RIGHT
      - Slide 3 = CENTER
      - Slide 4 = LEFT  (already left in initial — precondition, not scored)
      - Slide 5 = RIGHT
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Sanity check: presentation must have at least 7 slides
    num_slides = len(prs.slides)
    if num_slides < 7:
        print(f"CRITICAL: Expected 7 slides, found {num_slides}. File may be corrupt.")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Loaded presentation with {num_slides} slides.")

    # Component 1: Slide 2 must be RIGHT-aligned (0.35 points)
    # In initial_env, slide 2 is LEFT-aligned. This check FAILS on initial and PASSES on golden.
    try:
        slide2 = prs.slides[1]  # 0-indexed
        all_right, right_count, total2 = all_paragraphs_have_alignment(slide2, PP_ALIGN.RIGHT)
        if all_right:
            print(f"PASS: Component 1 — Slide 2 all {total2} paragraphs are RIGHT-aligned (0.35 pts)")
            total_score += 0.35
        else:
            # Check what alignment is actually set
            aligns = get_dominant_alignment(slide2)
            print(f"FAIL: Component 1 — Slide 2: expected all RIGHT, found {right_count}/{total2} RIGHT. Alignments: {aligns}")
    except Exception as e:
        print(f"ERROR: Component 1 (Slide 2) — {e}")

    # Component 2: Slide 3 must be CENTER-aligned (0.35 points)
    # In initial_env, slide 3 is LEFT-aligned. This check FAILS on initial and PASSES on golden.
    try:
        slide3 = prs.slides[2]  # 0-indexed
        all_center, center_count, total3 = all_paragraphs_have_alignment(slide3, PP_ALIGN.CENTER)
        if all_center:
            print(f"PASS: Component 2 — Slide 3 all {total3} paragraphs are CENTER-aligned (0.35 pts)")
            total_score += 0.35
        else:
            aligns = get_dominant_alignment(slide3)
            print(f"FAIL: Component 2 — Slide 3: expected all CENTER, found {center_count}/{total3} CENTER. Alignments: {aligns}")
    except Exception as e:
        print(f"ERROR: Component 2 (Slide 3) — {e}")

    # Component 3: Slide 5 must be RIGHT-aligned (0.30 points)
    # In initial_env, slide 5 is LEFT-aligned. This check FAILS on initial and PASSES on golden.
    try:
        slide5 = prs.slides[4]  # 0-indexed
        all_right5, right5_count, total5 = all_paragraphs_have_alignment(slide5, PP_ALIGN.RIGHT)
        if all_right5:
            print(f"PASS: Component 3 — Slide 5 all {total5} paragraphs are RIGHT-aligned (0.30 pts)")
            total_score += 0.30
        else:
            aligns = get_dominant_alignment(slide5)
            print(f"FAIL: Component 3 — Slide 5: expected all RIGHT, found {right5_count}/{total5} RIGHT. Alignments: {aligns}")
    except Exception as e:
        print(f"ERROR: Component 3 (Slide 5) — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in a given env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
