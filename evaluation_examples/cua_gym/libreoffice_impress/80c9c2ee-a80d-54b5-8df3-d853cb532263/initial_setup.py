"""
Initial Setup: 7-slide sales pitch deck with all slides (2-5) left-aligned text.
Task ID: osworld_impress_per_slide_alignment_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_per_slide_alignment_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_paragraph_alignment(text_frame, alignment):
    """Set all paragraphs in a text frame to the given alignment."""
    for para in text_frame.paragraphs:
        para.alignment = alignment


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "Nexus Solutions"
    sub1 = slide1.placeholders[1]
    sub1.text = "Transforming Business Through Innovation\n2025 Strategic Sales Pitch"
    for tf in [title1.text_frame, sub1.text_frame]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 2: Company Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Who We Are"
    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "Founded in 2012 with a mission to empower enterprises"
    p2 = tf2.add_paragraph()
    p2.text = "Over 500 global clients across 35 countries"
    p3 = tf2.add_paragraph()
    p3.text = "Annual revenue exceeding $120M in 2024"
    p4 = tf2.add_paragraph()
    p4.text = "Industry-leading Net Promoter Score of 72"
    p5 = tf2.add_paragraph()
    p5.text = "Certified ISO 27001 and SOC 2 Type II compliant"
    for tf in [title2.text_frame, tf2]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 3: Market Opportunity ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Market Opportunity"
    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Total addressable market valued at $4.2 billion by 2027"
    p3a = tf3.add_paragraph()
    p3a.text = "Digital transformation spending growing at 18.5% CAGR"
    p3b = tf3.add_paragraph()
    p3b.text = "Enterprise software market expanding across APAC and EMEA"
    p3c = tf3.add_paragraph()
    p3c.text = "60% of Fortune 500 companies planning cloud migration in 2025"
    p3d = tf3.add_paragraph()
    p3d.text = "Regulatory mandates driving compliance software demand"
    for tf in [title3.text_frame, tf3]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 4: Our Solution ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Our Solution"
    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "NexusPlatform — unified AI-driven operations suite"
    p4a = tf4.add_paragraph()
    p4a.text = "Real-time analytics with sub-50ms query response"
    p4b = tf4.add_paragraph()
    p4b.text = "Seamless ERP and CRM integrations via open APIs"
    p4c = tf4.add_paragraph()
    p4c.text = "Automated compliance reporting saving 300+ hours/year"
    p4d = tf4.add_paragraph()
    p4d.text = "Dedicated 24/7 customer success team"
    for tf in [title4.text_frame, tf4]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 5: Case Studies ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Customer Success Stories"
    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Meridian Healthcare: 40% reduction in operational costs"
    p5a = tf5.add_paragraph()
    p5a.text = "Apex Manufacturing: On-time delivery improved from 72% to 96%"
    p5b = tf5.add_paragraph()
    p5b.text = "Crestwood Financial: Regulatory audit time cut by 65%"
    p5c = tf5.add_paragraph()
    p5c.text = "GlobalTrade Corp: $8.3M saved in supply chain optimization"
    p5d = tf5.add_paragraph()
    p5d.text = "Deployed at scale in under 6 weeks average"
    for tf in [title5.text_frame, tf5]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 6: Pricing & Packages ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Flexible Pricing"
    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "Starter: $1,200/month — up to 25 users, core modules"
    p6a = tf6.add_paragraph()
    p6a.text = "Professional: $3,500/month — up to 100 users, full suite"
    p6b = tf6.add_paragraph()
    p6b.text = "Enterprise: Custom pricing — unlimited users, dedicated infra"
    p6c = tf6.add_paragraph()
    p6c.text = "All plans include onboarding, training, and SLA guarantee"
    p6d = tf6.add_paragraph()
    p6d.text = "30-day free trial available for qualified enterprises"
    for tf in [title6.text_frame, tf6]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    # --- Slide 7: Call to Action ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[0])
    title7 = slide7.shapes.title
    title7.text = "Let's Build the Future Together"
    sub7 = slide7.placeholders[1]
    sub7.text = "Contact us: sales@nexussolutions.com\n+1 (800) 555-0192\nwww.nexussolutions.com"
    for tf in [title7.text_frame, sub7.text_frame]:
        set_paragraph_alignment(tf, PP_ALIGN.LEFT)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
