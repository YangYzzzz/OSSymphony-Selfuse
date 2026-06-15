"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 191 is sitting on a deep navy background (#003366). Could you recolor the title text to its exact complementary shade, #FFCC00, so the heading really stands out?
Generated: 2025-09-10 18:35:40
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation
from pptx.dml.color import RGBColor

"""
Reward script for verifying the task:
  "Slide 191 is sitting on a deep navy background (#003366). Could you recolor the title text to its exact complementary shade, #FFCC00, so the heading really stands out?"

Scoring logic (progressive):
  • 0.4 points ‒ At least one slide with the exact navy background (#003366) exists.
  • 0.6 points ‒ Proportion of title-text runs on those slides that are exactly #FFCC00.
      - If every title-run is #FFCC00 → +0.6 (full).
      - Partial credit given proportionally.
  • Score is capped at 1.0 and rounded to two decimals.

Verification details:
  1. Load the PPTX safely (no points for merely loading).
  2. Detect slides whose background fill is solid #003366.
  3. Inspect each such slide’s title shape; examine every run’s font color.
  4. Compute accuracy and award proportional credit.

Prints detailed diagnostics and finally outputs:  "REWARD: X.X".
"""

def verify_title_color(file_path: str) -> float:
    navy_hex = "003366"  # Deep navy background colour
    gold_hex = "FFCC00"  # Complementary gold colour for title text

    # ---------- 0) File existence ----------
    if not os.path.exists(file_path):
        print("✗ File not found:", file_path)
        return 0.0

    # ---------- 1) Load presentation ----------
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("✗ Error loading presentation:", e)
        return 0.0

    score = 0.0

    # ---------- 2) Identify slides with correct background (#003366) ----------
    slides_with_navy_bg = []
    for idx, slide in enumerate(prs.slides):
        try:
            bg_fill = slide.background.fill
            if bg_fill.type == 1:  # SOLID
                rgb = bg_fill.fore_color.rgb  # may return RGBColor or None
                if rgb and str(rgb).upper() == navy_hex:
                    slides_with_navy_bg.append((idx, slide))
                    print(f"✓ Slide {idx + 1} has background #{navy_hex}")
        except Exception as e:
            print(f"   ! Could not evaluate background on slide {idx + 1}: {e}")

    if not slides_with_navy_bg:
        print("✗ No slide with background #003366 found → 0 points for background requirement")
        return 0.0  # Without the correct background the task is unfulfilled

    # Award 0.4 for finding the correct background slide(s)
    score += 0.4

    # ---------- 3) Verify title text colour (#FFCC00) ----------
    total_runs = 0
    correct_runs = 0

    for idx, slide in slides_with_navy_bg:
        title_shape = slide.shapes.title
        if title_shape is None or not title_shape.has_text_frame:
            print(f"✗ Slide {idx + 1} missing usable title shape")
            continue

        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                total_runs += 1
                run_rgb = run.font.color.rgb  # may be None
                if run_rgb and str(run_rgb).upper() == gold_hex:
                    correct_runs += 1
                else:
                    print(
                        f"✗ Slide {idx + 1} run colour {run_rgb} ≠ #{gold_hex}")

    if total_runs == 0:
        print("✗ No title runs found on slides with the navy background")
        return score  # Only the 0.4 for background

    # Compute proportion of correctly coloured runs
    accuracy = correct_runs / total_runs
    print(f"Title colour accuracy: {correct_runs}/{total_runs} = {accuracy:.2%}")

    # Award up to 0.6 based on accuracy
    score += 0.6 * accuracy

    final_score = round(min(score, 1.0), 2)
    print(f"Final task score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_191_is_sitting_on_a_deep_navy_background_003366_could_you_recolor_the_title_text_to_its_exact__golden.pptx"
    reward = verify_title_color(FILE_PATH)
    print(f"REWARD: {reward}")

