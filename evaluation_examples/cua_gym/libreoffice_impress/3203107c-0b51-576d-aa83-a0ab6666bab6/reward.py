"""
Reward Script: Create a simple agenda slide deck for team standup
Task ID: impress_wf_008
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.15): File exists with exactly 5 slides
  - Component 2 (0.20): Slide 1 title "Team Standup" with date in subtitle
  - Component 3 (0.20): Slides 2-5 titled "Topic 1" through "Topic 4"
  - Component 4 (0.25): Slides 2-5 have numbered list "1. Update / 2. Blockers / 3. Next steps"
  - Component 5 (0.20): All 5 slides have a blue (#0066CC) decorative shape near bottom
"""

import os
import re

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_008'


def get_slide_title(slide):
    """Extract title text from a slide."""
    if slide.shapes.title is not None:
        return slide.shapes.title.text.strip()
    # Fallback: look for shape named 'Title*'
    for shape in slide.shapes:
        if shape.has_text_frame and 'title' in shape.name.lower():
            return shape.text_frame.text.strip()
    return ""


def get_content_text(slide):
    """Get non-title text content from slide as list of paragraph texts."""
    paragraphs = []
    for shape in slide.shapes:
        if shape.has_text_frame and 'title' not in shape.name.lower():
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    paragraphs.append(text)
    return paragraphs


def has_blue_line_shape(slide, target_rgb='0066CC'):
    """
    Check if slide has a shape with blue fill near the bottom that acts as decorative line.
    Looks for a thin shape (height < 100000 EMU ~ 0.11 in) in the lower portion
    with solid fill color matching target_rgb.
    """
    slide_height = 6858000  # default slide height in EMU (7.5 in)
    for shape in slide.shapes:
        # Check if shape is thin (line-like) and near bottom
        if shape.height < 200000 and shape.top > slide_height * 0.7:
            # Check fill color
            try:
                fill = shape.fill
                if fill.type is not None and fill.type == 1:  # SOLID
                    color = str(fill.fore_color.rgb).upper()
                    if color == target_rgb.upper():
                        return True
            except Exception:
                pass
            # Also check via XML for srgbClr
            try:
                xml_str = shape.element.xml
                colors = re.findall(r'srgbClr val="([A-Fa-f0-9]{6})"', xml_str)
                if target_rgb.upper() in [c.upper() for c in colors]:
                    if shape.height < 200000 and shape.top > slide_height * 0.7:
                        return True
            except Exception:
                pass
    return False


def _check_slide_has_date(slide):
    """Check if slide has a date-like text in non-title shapes."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_title_shape = (shape == slide.shapes.title) or \
                         (shape.name.lower().startswith('title') and 'sub' not in shape.name.lower())
        if not is_title_shape:
            text = shape.text_frame.text.strip()
            if re.search(r'\d{1,4}[-/.,\s]\d{1,2}[-/.,\s]\d{1,4}', text) or \
               re.search(r'(january|february|march|april|may|june|july|august|september|october|november|december)\s+\d{1,2}', text, re.IGNORECASE) or \
               re.search(r'\d{1,2}\s+(january|february|march|april|may|june|july|august|september|october|november|december)', text, re.IGNORECASE) or \
               re.search(r'(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\w*\s+\d{1,2}', text, re.IGNORECASE):
                return True
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: File has exactly 5 slides (0.15 points)
    try:
        if num_slides == 5:
            print(f"PASS: Component 1 — File has 5 slides (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 — Expected 5 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 1 title is "Team Standup" with date in subtitle (0.20 points)
    try:
        if num_slides >= 1:
            slide1 = prs.slides[0]
            title_text = get_slide_title(slide1)
            # Check title
            title_match = title_text.lower() == 'team standup'
            # Check subtitle/date - look for date-like text in non-title shapes
            has_date = _check_slide_has_date(slide1)

            if title_match and has_date:
                print(f"PASS: Component 2 — Slide 1 title='Team Standup' with date (0.20 pts)")
                total_score += 0.20
            elif title_match:
                print(f"PARTIAL: Component 2 — Title correct but no date found (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 — Title='{title_text}', date found={has_date}")
        else:
            print(f"FAIL: Component 2 — No slides present")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slides 2-5 titled "Topic 1" through "Topic 4" (0.20 points)
    try:
        expected_titles = ['Topic 1', 'Topic 2', 'Topic 3', 'Topic 4']
        correct_titles = 0
        for idx in range(1, min(5, num_slides)):
            slide = prs.slides[idx]
            title = get_slide_title(slide)
            expected = expected_titles[idx - 1]
            if title.lower() == expected.lower():
                correct_titles += 1
            else:
                print(f"  Slide {idx+1}: expected '{expected}', found '{title}'")

        if correct_titles == 4:
            print(f"PASS: Component 3 — All 4 topic slides correctly titled (0.20 pts)")
            total_score += 0.20
        elif correct_titles > 0:
            partial = round(0.20 * correct_titles / 4, 2)
            print(f"PARTIAL: Component 3 — {correct_titles}/4 titles correct ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 — No topic titles match")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slides 2-5 have numbered list content (0.25 points)
    # Each slide should have "1. Update", "2. Blockers", "3. Next steps"
    try:
        expected_items = ['1. update', '2. blockers', '3. next steps']
        slides_with_list = 0
        for idx in range(1, min(5, num_slides)):
            slide = prs.slides[idx]
            content = get_content_text(slide)
            content_lower = [c.lower().strip() for c in content]
            # Check if all 3 expected items are present
            all_found = all(
                any(exp in c for c in content_lower)
                for exp in expected_items
            )
            if all_found:
                slides_with_list += 1
            else:
                print(f"  Slide {idx+1} content: {content}")

        if slides_with_list == 4:
            print(f"PASS: Component 4 — All 4 topic slides have numbered list (0.25 pts)")
            total_score += 0.25
        elif slides_with_list > 0:
            partial = round(0.25 * slides_with_list / 4, 2)
            print(f"PARTIAL: Component 4 — {slides_with_list}/4 slides have list ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 — No slides have the expected numbered list")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: All slides have blue (#0066CC) decorative line near bottom (0.20 points)
    try:
        slides_with_line = 0
        for idx in range(min(5, num_slides)):
            slide = prs.slides[idx]
            if has_blue_line_shape(slide):
                slides_with_line += 1
            else:
                print(f"  Slide {idx+1}: no blue decorative line found")

        target_count = min(5, num_slides)
        if slides_with_line == target_count and target_count == 5:
            print(f"PASS: Component 5 — All 5 slides have blue line (0.20 pts)")
            total_score += 0.20
        elif slides_with_line > 0:
            partial = round(0.20 * slides_with_line / 5, 2)
            print(f"PARTIAL: Component 5 — {slides_with_line}/5 slides have blue line ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 5 — No slides have blue decorative line")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = min(round(total_score, 2), 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/Desktop/Standup.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
