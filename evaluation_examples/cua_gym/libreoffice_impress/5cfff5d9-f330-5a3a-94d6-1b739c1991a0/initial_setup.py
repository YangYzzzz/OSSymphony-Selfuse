#!/usr/bin/env python3
"""
initial_setup.py — Create a blank presentation for impress_gf4_022.
Creates a blank .pptx with default slide size (25.4 x 19.05 cm), one blank slide.
Opens it in LibreOffice Impress.
"""

import os
import subprocess
import shlex
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_gf4_022'

# Install dependency
subprocess.run(['pip3', 'install', 'python-pptx'], capture_output=True)

from pptx import Presentation
from pptx.util import Emu

# Create blank presentation with default slide size
prs = Presentation()

# Default slide size: 25.4 cm x 19.05 cm (already the default in python-pptx)
# 25.4 cm = 9144000 EMU, 19.05 cm = 6858000 EMU
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(6858000)

# Add one blank slide (layout 6 = Blank)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Save
filepath = os.path.join(WORKDIR, f'{TASK_ID}_initial.pptx')
prs.save(filepath)
print(f"Saved initial file to {filepath}")

# Verify file exists
assert os.path.exists(filepath), f"File not found: {filepath}"
print(f"File size: {os.path.getsize(filepath)} bytes")

# Open in LibreOffice Impress
env = os.environ.copy()
env["DISPLAY"] = ":0"
subprocess.Popen(
    shlex.split(f'libreoffice --impress "{filepath}"'),
    stdout=subprocess.DEVNULL,
    stderr=subprocess.DEVNULL,
    env=env,
)
time.sleep(2)
print("LibreOffice Impress launched.")
