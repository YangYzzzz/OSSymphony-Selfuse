#!/usr/bin/env python3
"""
Reward script for impress_gf4_022: Scientific conference poster at A0 size.
Verifies slide size, header, title, authors, 3-column layout, flowchart,
bar chart, conclusions bullets, and color scheme.
"""

import glob
import os
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def find_pptx():
    """Find the task pptx file on this VM."""
    patterns = [
        "/home/user/impress_gf4_022_golden.pptx",
        "/home/user/impress_gf4_022_initial.pptx",
        "/home/user/impress_gf4_022*.pptx",
    ]
    for pat in patterns:
        matches = glob.glob(pat)
        if matches:
            return matches[0]
    # Fallback: any pptx in /home/user
    matches = glob.glob("/home/user/*.pptx")
    if matches:
        return matches[0]
    return None


def approx_equal(a, b, tol=0.03):
    """Check approximate equality with relative tolerance."""
    if a == b:
        return True
    if a == 0 or b == 0:
        return abs(a - b) < 100000  # small absolute tolerance for zero
    return abs(a - b) / max(abs(a), abs(b)) <= tol


def get_shape_fill_color(shape):
    """Get the solid fill RGB color of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_run_color(run):
    """Get the RGB color string of a run, or None."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def check_bullets_in_pptx(pptx_path):
    """Check for bullet characters in conclusions text box via XML."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    bullet_count = 0
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            with zf.open("ppt/slides/slide1.xml") as f:
                root = ET.parse(f).getroot()
                for p in root.findall(".//{%s}p" % ns_a):
                    pPr = p.find("{%s}pPr" % ns_a)
                    if pPr is not None:
                        bc = pPr.find("{%s}buChar" % ns_a)
                        if bc is not None and bc.get("char"):
                            txt = "".join(
                                t.text or "" for t in p.findall(".//{%s}t" % ns_a)
                            )
                            if txt.strip():
                                bullet_count += 1
    except Exception:
        pass
    return bullet_count


def evaluate(pptx_path):
    score = 0.0
    details = []

    prs = Presentation(pptx_path)
    if len(prs.slides) == 0:
        print("No slides found.")
        return 0.0

    slide = prs.slides[0]
    shapes = list(slide.shapes)

    # =========================================================
    # 1. Slide size: A0 landscape (118.9 x 84.1 cm)
    #    118.9 cm = 42,804,000 EMU, 84.1 cm = 30,276,000 EMU
    # =========================================================
    w_ok = approx_equal(prs.slide_width, 42804000, tol=0.02)
    h_ok = approx_equal(prs.slide_height, 30276000, tol=0.02)
    landscape_ok = prs.slide_width > prs.slide_height
    if w_ok and h_ok and landscape_ok:
        score += 0.15
        details.append("PASS slide_size (0.15)")
    else:
        details.append(
            f"FAIL slide_size: w={prs.slide_width} (expect ~42804000), "
            f"h={prs.slide_height} (expect ~30276000), landscape={landscape_ok}"
        )

    # =========================================================
    # 2. Dark blue header rectangle (#003366), full-width at top
    # =========================================================
    header_rect = None
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or sh.shape_type == 1:
            color = get_shape_fill_color(sh)
            if color and color.upper() == "003366":
                # Full-width: shape width should be ~slide width
                if approx_equal(sh.width, prs.slide_width, tol=0.05):
                    # At top: top should be near 0
                    if sh.top < 1000000:  # within ~2.8cm of top
                        header_rect = sh
                        break

    if header_rect is not None:
        score += 0.15
        details.append("PASS header_rect (0.15)")
    else:
        details.append("FAIL header_rect: no full-width #003366 rectangle at top")

    # =========================================================
    # 3. Title text: 48pt white bold
    # =========================================================
    title_found = False
    for sh in shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    rcolor = get_run_color(r)
                    is_white = rcolor and rcolor.upper() == "FFFFFF"
                    is_bold = r.font.bold is True
                    is_48pt = r.font.size is not None and approx_equal(
                        r.font.size, Pt(48), tol=0.05
                    )
                    if is_white and is_bold and is_48pt and len(r.text.strip()) > 5:
                        title_found = True
                        break
                if title_found:
                    break
        if title_found:
            break

    if title_found:
        score += 0.10
        details.append("PASS title_text (0.10)")
    else:
        details.append("FAIL title_text: no 48pt white bold text found")

    # =========================================================
    # 4. Author names: 24pt white text
    # =========================================================
    author_found = False
    for sh in shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    rcolor = get_run_color(r)
                    is_white = rcolor and rcolor.upper() == "FFFFFF"
                    is_24pt = r.font.size is not None and approx_equal(
                        r.font.size, Pt(24), tol=0.05
                    )
                    if is_white and is_24pt and len(r.text.strip()) > 5:
                        author_found = True
                        break
                if author_found:
                    break
        if author_found:
            break

    if author_found:
        score += 0.10
        details.append("PASS author_text (0.10)")
    else:
        details.append("FAIL author_text: no 24pt white text found")

    # =========================================================
    # 5. Three column rectangles below header
    # =========================================================
    col_rects = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or sh.shape_type == 1:
            color = get_shape_fill_color(sh)
            # Column rects should be below the header (top > 5,000,000 EMU)
            # and not be the header itself
            if sh.top > 4000000 and sh.height > 10000000:
                col_rects.append(sh)

    # Sort by left position
    col_rects.sort(key=lambda s: s.left)

    if len(col_rects) >= 3:
        # Check roughly equal widths
        widths = [s.width for s in col_rects[:3]]
        widths_similar = all(
            approx_equal(widths[0], w, tol=0.1) for w in widths[1:]
        )
        if widths_similar:
            score += 0.15
            details.append("PASS column_rects (0.15)")
        else:
            details.append(
                f"FAIL column_rects: widths not equal: {widths}"
            )
    else:
        details.append(
            f"FAIL column_rects: found {len(col_rects)} tall rects below header (need 3)"
        )

    # =========================================================
    # 6. Flowchart shapes in column 1: at least 5 shapes
    #    (rounded rects + arrows in the left third)
    # =========================================================
    slide_third = prs.slide_width / 3
    flowchart_shapes = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or sh.shape_type == 1:
            # In left third of slide and below header
            if sh.left < slide_third and sh.top > 5000000:
                # Exclude the column background rect itself
                if sh.height < 10000000:
                    flowchart_shapes.append(sh)

    if len(flowchart_shapes) >= 5:
        score += 0.10
        details.append(f"PASS flowchart_shapes (0.10): {len(flowchart_shapes)} shapes")
    else:
        details.append(
            f"FAIL flowchart_shapes: found {len(flowchart_shapes)} shapes in column 1 (need >=5)"
        )

    # =========================================================
    # 7. Bar chart in column 2: embedded chart object
    # =========================================================
    chart_found = False
    for sh in shapes:
        if hasattr(sh, "has_chart") and sh.has_chart:
            # Should be in middle third
            if sh.left >= slide_third * 0.8:
                chart_found = True
                break

    if chart_found:
        score += 0.10
        details.append("PASS bar_chart (0.10)")
    else:
        details.append("FAIL bar_chart: no embedded chart found in column 2 area")

    # =========================================================
    # 8. Conclusions text with bullet icons in column 3
    # =========================================================
    bullet_count = check_bullets_in_pptx(pptx_path)
    # Also check there's text in the right third
    conclusions_text_found = False
    for sh in shapes:
        if sh.has_text_frame and sh.left > slide_third * 1.5:
            text = sh.text_frame.text.strip()
            if len(text) > 50:
                conclusions_text_found = True
                break

    if bullet_count >= 3 and conclusions_text_found:
        score += 0.10
        details.append(f"PASS conclusions_bullets (0.10): {bullet_count} bullets")
    else:
        details.append(
            f"FAIL conclusions_bullets: bullets={bullet_count}, text_found={conclusions_text_found}"
        )

    # =========================================================
    # 9. Color scheme consistency: #003366 and #F0F4FF used
    # =========================================================
    dark_blue_count = 0
    light_grey_count = 0
    for sh in shapes:
        color = get_shape_fill_color(sh)
        if color:
            if color.upper() == "003366":
                dark_blue_count += 1
            if color.upper() == "F0F4FF":
                light_grey_count += 1

    if dark_blue_count >= 2 and light_grey_count >= 1:
        score += 0.05
        details.append(
            f"PASS color_scheme (0.05): dark_blue={dark_blue_count}, light_grey={light_grey_count}"
        )
    else:
        details.append(
            f"FAIL color_scheme: dark_blue={dark_blue_count}, light_grey={light_grey_count}"
        )

    for d in details:
        print(d)

    return round(score, 2)


def main():
    pptx_path = find_pptx()
    if pptx_path is None:
        print("ERROR: No pptx file found in /home/user/")
        print("REWARD: 0.0")
        return

    print(f"Evaluating: {pptx_path}")
    score = evaluate(pptx_path)
    print(f"REWARD: {score}")


if __name__ == "__main__":
    main()
