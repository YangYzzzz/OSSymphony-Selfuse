"""
Reward Script: Change slide 2 title to 'Overview of Key Findings' and match color to slide 1's #0D1B2A
Task ID: osworld_impress_title_color_match_010
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5 pts): Slide 2 title text == 'Overview of Key Findings'
  Component 2 (0.5 pts): Slide 2 title color == #0D1B2A
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_010'

EXPECTED_TITLE_TEXT = 'Overview of Key Findings'
EXPECTED_COLOR = '0D1B2A'  # dark navy — matches slide 1 title color


def get_title_shape(slide):
    """Return the title placeholder shape from a slide, or None."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            # idx 0 is title/center-title
            if ph.idx == 0 and shape.has_text_frame:
                return shape
    return None


def get_run_color_hex(run):
    """Return the hex color string of a run's font, or None if not set."""
    try:
        if run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load the presentation — gate: if we cannot load, score 0.0
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Gate: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"CRITICAL: Expected at least 2 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed → slide 2

    title_shape = get_title_shape(slide2)
    if title_shape is None:
        print("CRITICAL: No title placeholder found on slide 2")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 title text == 'Overview of Key Findings' (0.5 points)
    try:
        actual_text = title_shape.text_frame.text.strip()
        if actual_text == EXPECTED_TITLE_TEXT:
            print(f"PASS: Component 1 — Slide 2 title text is '{actual_text}' (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — Expected title '{EXPECTED_TITLE_TEXT}', found '{actual_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 — Could not check slide 2 title text: {e}")

    # Component 2: Slide 2 title color == #0D1B2A (0.5 points)
    # Check all non-empty runs in the title placeholder
    try:
        paragraphs = title_shape.text_frame.paragraphs
        runs_with_text = [
            run
            for para in paragraphs
            for run in para.runs
            if (run.text or '').strip()
        ]

        if not runs_with_text:
            print(f"FAIL: Component 2 — No text runs found in slide 2 title")
        else:
            # Collect mismatched colors
            mismatches = [
                (run.text, get_run_color_hex(run))
                for run in runs_with_text
                if get_run_color_hex(run) != EXPECTED_COLOR
            ]

            if mismatches:
                for text, color_hex in mismatches:
                    print(f"FAIL: Component 2 — Run '{text}' has color {color_hex!r}, expected '{EXPECTED_COLOR}'")
            else:
                print(f"PASS: Component 2 — Slide 2 title color is #{EXPECTED_COLOR} (dark navy, matches slide 1) (0.5 pts)")
                if not mismatches:
                    total_score += 0.5
    except Exception as e:
        print(f"ERROR: Component 2 — Could not check slide 2 title color: {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM environment
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
