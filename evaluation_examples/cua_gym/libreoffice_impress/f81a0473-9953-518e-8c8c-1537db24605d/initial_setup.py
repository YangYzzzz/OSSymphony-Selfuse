"""
Initial Setup: 5-slide report deck where slide 2 title reads 'Overview' in black
Task ID: osworld_impress_title_color_match_010
Domain: libreoffice_impress

Creates a pre-task presentation where:
  - Slide 1 has a title in dark navy #0D1B2A
  - Slide 2 has a title 'Overview' in black (not #0D1B2A, not matching slide 1)
The agent must rename slide 2 title to 'Overview of Key Findings' and apply #0D1B2A.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_color(slide, color_rgb):
    """Set the color of all runs in the title placeholder."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                return
    # fallback: find title by name or first title-like shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            if 'title' in shape.name.lower() or 'Title' in shape.name:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                return


def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(0x0D, 0x1B, 0x2A)
    black = RGBColor(0x00, 0x00, 0x00)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    dark_gray = RGBColor(0x33, 0x33, 0x33)
    accent_blue = RGBColor(0x1F, 0x6F, 0xEB)

    # ---- Slide 1: Title slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Q4 Performance Report"
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    slide1.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
    try:
        slide1.placeholders[1].text = "Annual Review — Fiscal Year 2024"
        slide1.placeholders[1].text_frame.paragraphs[0].runs[0].font.size = Pt(20)
        slide1.placeholders[1].text_frame.paragraphs[0].runs[0].font.color.rgb = dark_gray
    except (KeyError, IndexError):
        pass

    # ---- Slide 2: 'Overview' in black (NOT navy) — agent must fix this ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Overview"
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    slide2.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = black
    try:
        content2 = slide2.placeholders[1]
        tf2 = content2.text_frame
        tf2.text = "This section summarizes the major highlights from Q4 2024."
        tf2.paragraphs[0].runs[0].font.size = Pt(18)
        tf2.paragraphs[0].runs[0].font.color.rgb = dark_gray
        p2b = tf2.add_paragraph()
        p2b.text = "Revenue grew by 12.4% compared to Q3 2024"
        p2b.runs[0].font.size = Pt(16)
        p2b.runs[0].font.color.rgb = dark_gray
        p2c = tf2.add_paragraph()
        p2c.text = "Customer retention reached an all-time high of 94.7%"
        p2c.runs[0].font.size = Pt(16)
        p2c.runs[0].font.color.rgb = dark_gray
        p2d = tf2.add_paragraph()
        p2d.text = "Operating costs reduced by 8.2% through process optimization"
        p2d.runs[0].font.size = Pt(16)
        p2d.runs[0].font.color.rgb = dark_gray
    except (KeyError, IndexError):
        pass

    # ---- Slide 3: Revenue Analysis ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Revenue Analysis"
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    slide3.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
    try:
        content3 = slide3.placeholders[1]
        tf3 = content3.text_frame
        tf3.text = "Q4 2024 Revenue Breakdown by Region"
        tf3.paragraphs[0].runs[0].font.size = Pt(18)
        tf3.paragraphs[0].runs[0].font.color.rgb = dark_gray
        region_data = [
            ("North America", "$3.2M", "+15.3%"),
            ("Europe", "$2.1M", "+9.8%"),
            ("Asia Pacific", "$1.8M", "+21.4%"),
            ("Latin America", "$0.6M", "+6.7%"),
        ]
        for region, revenue, growth in region_data:
            p = tf3.add_paragraph()
            p.text = f"{region}: {revenue} ({growth})"
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = dark_gray
    except (KeyError, IndexError):
        pass

    # ---- Slide 4: Customer Insights ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Customer Insights"
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    slide4.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
    try:
        content4 = slide4.placeholders[1]
        tf4 = content4.text_frame
        tf4.text = "Key metrics from our customer satisfaction survey (n=2,847)"
        tf4.paragraphs[0].runs[0].font.size = Pt(18)
        tf4.paragraphs[0].runs[0].font.color.rgb = dark_gray
        metrics = [
            "Net Promoter Score (NPS): 68 (industry avg: 42)",
            "Average response time: 2.4 hours (down from 4.1 hrs)",
            "First-contact resolution rate: 87.3%",
            "Premium tier adoption: 34.5% of active accounts",
        ]
        for m in metrics:
            p = tf4.add_paragraph()
            p.text = m
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = dark_gray
    except (KeyError, IndexError):
        pass

    # ---- Slide 5: Outlook & Next Steps ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Outlook & Next Steps"
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    slide5.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = navy
    try:
        content5 = slide5.placeholders[1]
        tf5 = content5.text_frame
        tf5.text = "Strategic priorities for Q1 2025:"
        tf5.paragraphs[0].runs[0].font.size = Pt(18)
        tf5.paragraphs[0].runs[0].font.color.rgb = dark_gray
        next_steps = [
            "Launch product line expansion in APAC markets (March 2025)",
            "Complete CRM system migration and staff training by end of Q1",
            "Achieve 95%+ customer retention target through proactive outreach",
            "Reduce average deal closure time from 32 to 22 days",
        ]
        for ns in next_steps:
            p = tf5.add_paragraph()
            p.text = ns
            p.runs[0].font.size = Pt(15)
            p.runs[0].font.color.rgb = dark_gray
    except (KeyError, IndexError):
        pass

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
