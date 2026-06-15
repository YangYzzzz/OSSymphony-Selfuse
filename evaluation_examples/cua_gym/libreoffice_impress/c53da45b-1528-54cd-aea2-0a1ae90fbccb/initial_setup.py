"""
Initial Setup: Photography portfolio presentation with title at top of slide 3
Task ID: osworld_impress_title_position_bottom_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_002'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # Use standard widescreen dimensions (13.33 x 7.5 inches)
    slide_width = prs.slide_width    # 9144000 EMU = 10 inches (default)
    slide_height = prs.slide_height  # 6858000 EMU = 7.5 inches (default)

    blank_layout = prs.slide_layouts[6]   # Blank layout
    title_only_layout = prs.slide_layouts[5]  # Title Only layout

    # ─── Slide 1: Cover ───
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Photography Portfolio"
    slide1.placeholders[1].text = "A Visual Journey Through Light and Perspective"
    # Dark background
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    # Style title
    title_run = slide1.shapes.title.text_frame.paragraphs[0].runs[0]
    title_run.font.color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    title_run.font.size = Pt(40)
    title_run.font.bold = True

    # ─── Slide 2: Nature Collection ───
    slide2 = prs.slides.add_slide(title_only_layout)
    # Title at top
    title2 = slide2.shapes.title
    title2.text = "Nature Collection"
    title2.left = Inches(0.5)
    title2.top = Inches(0.3)
    title2.width = Inches(9.0)
    title2.height = Inches(0.8)
    t2_run = title2.text_frame.paragraphs[0].runs[0]
    t2_run.font.size = Pt(28)
    t2_run.font.bold = True
    t2_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Dark background
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    # Description text box
    txBox2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.0))
    tf2 = txBox2.text_frame
    tf2.text = "Landscape and wildlife photography from national parks and wilderness areas"
    tf2.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    tf2.paragraphs[0].runs[0].font.size = Pt(16)

    # ─── Slide 3: Cityscape Collection — TITLE AT TOP ───
    slide3 = prs.slides.add_slide(title_only_layout)
    # Title positioned at the TOP of the slide (initial state)
    title3 = slide3.shapes.title
    title3.text = "Cityscape Collection"
    title3.left = Inches(0.5)
    title3.top = Inches(0.3)        # TOP position - this is what agent must move
    title3.width = Inches(9.0)
    title3.height = Inches(0.75)
    t3_run = title3.text_frame.paragraphs[0].runs[0]
    t3_run.font.size = Pt(28)
    t3_run.font.bold = True
    t3_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    # Dark background
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x17, 0x20, 0x2A)
    # Large image placeholder (represented as a rectangle for initial state)
    # Image area fills most of slide below the title
    img_box3 = slide3.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.5))
    img_tf3 = img_box3.text_frame
    img_tf3.text = "[City Skyline Photography]"
    img_tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
    img_tf3.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x55, 0x66, 0x77)
    img_tf3.paragraphs[0].runs[0].font.size = Pt(20)
    # Style the image area background
    img_box3.fill.solid()
    img_box3.fill.fore_color.rgb = RGBColor(0x22, 0x33, 0x44)

    # ─── Slide 4: Portrait Collection ───
    slide4 = prs.slides.add_slide(title_only_layout)
    title4 = slide4.shapes.title
    title4.text = "Portrait Collection"
    title4.left = Inches(0.5)
    title4.top = Inches(0.3)
    title4.width = Inches(9.0)
    title4.height = Inches(0.8)
    t4_run = title4.text_frame.paragraphs[0].runs[0]
    t4_run.font.size = Pt(28)
    t4_run.font.bold = True
    t4_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x2E, 0x1A, 0x2E)
    txBox4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.0))
    tf4 = txBox4.text_frame
    tf4.text = "Candid and studio portraits capturing authentic human expressions"
    tf4.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    tf4.paragraphs[0].runs[0].font.size = Pt(16)

    # ─── Slide 5: Abstract Collection ───
    slide5 = prs.slides.add_slide(title_only_layout)
    title5 = slide5.shapes.title
    title5.text = "Abstract Collection"
    title5.left = Inches(0.5)
    title5.top = Inches(0.3)
    title5.width = Inches(9.0)
    title5.height = Inches(0.8)
    t5_run = title5.text_frame.paragraphs[0].runs[0]
    t5_run.font.size = Pt(28)
    t5_run.font.bold = True
    t5_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x1A)
    txBox5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.0))
    tf5 = txBox5.text_frame
    tf5.text = "Experimental compositions exploring form, color, and texture"
    tf5.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    tf5.paragraphs[0].runs[0].font.size = Pt(16)

    # ─── Slide 6: Contact & Booking ───
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Contact & Booking"
    slide6.placeholders[1].text = (
        "Elena Vasquez Photography\n"
        "elena@vasquezphoto.com\n"
        "+1 (555) 847-2391\n"
        "www.vasquezphoto.com\n"
        "Available for commercial and editorial assignments worldwide"
    )
    fill6 = slide6.background.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    title6_run = slide6.shapes.title.text_frame.paragraphs[0].runs[0]
    title6_run.font.color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    title6_run.font.size = Pt(32)
    title6_run.font.bold = True

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
