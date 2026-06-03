"""
Reward Script: Reposition the title on slide 3 so it appears at the bottom of the slide instead of the top.
Task ID: osworld_impress_title_position_bottom_002
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): Title 'Cityscape Collection' on slide 3 is moved to the lower half of the slide (top > slide_height/2)
  Component 2 (0.5): Title is fully in the bottom portion — bottom edge (top + height) >= 85% of slide_height
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_position_bottom_002'

def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Task: Move title on slide 3 from top to bottom of slide.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 3 slides
    try:
        num_slides = len(prs.slides)
        if num_slides < 3:
            print(f"PRECONDITION FAIL: Expected at least 3 slides, found {num_slides}")
            print("REWARD: 0.0")
            return 0.0
        print(f"OK: Presentation has {num_slides} slides")
    except Exception as e:
        print(f"CRITICAL: Cannot count slides: {e}")
        print("REWARD: 0.0")
        return 0.0

    slide_height = prs.slide_height  # in EMU
    slide = prs.slides[2]  # slide 3, 0-indexed

    # Find the title shape on slide 3 (shape named 'Title 1' or placeholder type title)
    title_shape = None
    try:
        for shape in slide.shapes:
            if shape.has_text_frame and 'Cityscape Collection' in shape.text_frame.text:
                title_shape = shape
                break
        if title_shape is None:
            # Also try by name
            for shape in slide.shapes:
                if shape.name and 'Title' in shape.name:
                    title_shape = shape
                    break
        if title_shape is None:
            print("FAIL: Could not find title shape with 'Cityscape Collection' on slide 3")
            print(f"REWARD: 0.0")
            return 0.0
        print(f"OK: Found title shape '{title_shape.name}' with text: '{title_shape.text_frame.text}'")
    except Exception as e:
        print(f"CRITICAL: Error finding title shape: {e}")
        print("REWARD: 0.0")
        return 0.0

    title_top = title_shape.top        # EMU from top of slide
    title_height = title_shape.height  # EMU height of the shape

    print(f"Title top: {title_top} EMU ({title_top/914400:.3f} inches)")
    print(f"Title height: {title_height} EMU ({title_height/914400:.3f} inches)")
    print(f"Slide height: {slide_height} EMU ({slide_height/914400:.3f} inches)")
    bottom_edge = title_top + title_height
    print(f"Title bottom edge: {bottom_edge} EMU ({bottom_edge/914400:.3f} inches)")

    # Component 1: Title top is in the lower half of the slide (top > slide_height / 2)
    # Initial state: top ≈ 274320 (0.300 inches, near top)
    # Golden state: top ≈ 5943600 (6.500 inches, near bottom)
    # This FAILS on initial (top = 274320 < slide_height/2 = ~3429000) and PASSES on golden
    try:
        lower_half_threshold = slide_height // 2
        if title_top > lower_half_threshold:
            print(f"PASS: Component 1 — title top ({title_top}) is in lower half of slide (threshold: {lower_half_threshold}) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 1 — title top ({title_top}) is NOT in lower half of slide (threshold: {lower_half_threshold})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Title is near the bottom — bottom edge (top + height) is >= 85% of slide_height
    # Initial state: bottom_edge = 274320 + 685800 = 960120 (~13% of slide), FAILS
    # Golden state: bottom_edge = 5943600 + 685800 = 6629400 (~97% of slide), PASSES
    try:
        threshold_85pct = int(slide_height * 0.85)
        if bottom_edge >= threshold_85pct:
            print(f"PASS: Component 2 — title bottom edge ({bottom_edge}) >= 85% slide_height ({threshold_85pct}) (0.5 pts)")
            total_score += 0.5
        else:
            print(f"FAIL: Component 2 — title bottom edge ({bottom_edge}) < 85% slide_height ({threshold_85pct})")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
