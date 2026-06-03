"""
Initial Setup: 6-slide process overview presentation with flat bullet list on slide 3
Task ID: osworld_impress_bullet_indent_remove_008
Domain: libreoffice_impress

Slide 3 has 6 flat (all top-level) bullet items — agent must restructure into hierarchy
and remove all bullet symbols.
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bullet_indent_remove_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_bullet_char(para, char='•'):
    """Set a bullet character on a paragraph using XML."""
    pPr = para._p.get_or_add_pPr()
    # Remove any existing bullet-related elements
    for tag in ['a:buNone', 'a:buChar', 'a:buAutoNum', 'a:buClrTx', 'a:buClr',
                'a:buSzTx', 'a:buSzPct', 'a:buSzPts', 'a:buFontTx', 'a:buFont']:
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)
    # Add buChar element
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', char)


def create_initial():
    prs = Presentation()
    # Use standard 10x7.5 inch widescreen layout
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layouts = prs.slide_layouts
    # Layout indices: 0=Title Slide, 1=Title+Content, 2=Title+2Content, etc.

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(slide_layouts[0])
    slide1.shapes.title.text = "Operational Excellence Framework"
    try:
        slide1.placeholders[1].text = "A 6-Phase Process Overview\nQ2 2025 Initiative"
    except (KeyError, IndexError):
        pass

    # ---- Slide 2: Introduction ----
    slide2 = prs.slides.add_slide(slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "This presentation outlines the six core phases of our operational excellence program."
    p2 = tf2.add_paragraph()
    p2.text = "Each phase builds on the previous to drive sustainable improvements."
    p2_2 = tf2.add_paragraph()
    p2_2.text = "Teams across all departments are expected to participate actively."

    # ---- Slide 3: Process Steps (6 flat bullet items, all top-level) ----
    slide3 = prs.slides.add_slide(slide_layouts[1])
    slide3.shapes.title.text = "Process Steps"
    tf3 = slide3.placeholders[1].text_frame
    tf3.word_wrap = True

    # Define 6 flat bullet items
    bullet_items = [
        "Define Scope and Objectives",
        "Conduct Stakeholder Analysis",
        "Map Current State Processes",
        "Identify Improvement Opportunities",
        "Develop Action Plans",
        "Implement and Monitor Results",
    ]

    # Set first paragraph text
    tf3.paragraphs[0].text = bullet_items[0]
    tf3.paragraphs[0].level = 0
    set_bullet_char(tf3.paragraphs[0], '•')

    # Add remaining bullet items as flat top-level paragraphs
    for item in bullet_items[1:]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0  # All at top level
        set_bullet_char(p, '•')

    # ---- Slide 4: Key Metrics ----
    slide4 = prs.slides.add_slide(slide_layouts[1])
    slide4.shapes.title.text = "Key Performance Metrics"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Process cycle time reduced by 23% in pilot groups"
    metrics = [
        "Customer satisfaction scores up 18 points",
        "Defect rate dropped from 4.2% to 1.1%",
        "On-time delivery improved to 96.7%",
        "Cost savings of $2.4M projected annually",
    ]
    for m in metrics:
        mp = tf4.add_paragraph()
        mp.text = m
        mp.level = 0
        set_bullet_char(mp, '•')

    # ---- Slide 5: Timeline ----
    slide5 = prs.slides.add_slide(slide_layouts[1])
    slide5.shapes.title.text = "Implementation Timeline"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Phase 1 — Planning (January – February 2025)"
    timeline = [
        "Phase 2 — Analysis (March – April 2025)",
        "Phase 3 — Design (May – June 2025)",
        "Phase 4 — Pilot (July – September 2025)",
        "Phase 5 — Rollout (October – December 2025)",
        "Phase 6 — Review (January 2026)",
    ]
    for t in timeline:
        tp = tf5.add_paragraph()
        tp.text = t
        tp.level = 0
        set_bullet_char(tp, '•')

    # ---- Slide 6: Conclusion ----
    slide6 = prs.slides.add_slide(slide_layouts[1])
    slide6.shapes.title.text = "Next Steps and Conclusion"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Commit to the process and allocate required resources"
    conclusions = [
        "Schedule kickoff meetings with department leads",
        "Establish reporting cadence and accountability",
        "Review progress quarterly against defined KPIs",
    ]
    for c in conclusions:
        cp = tf6.add_paragraph()
        cp.text = c
        cp.level = 0
        set_bullet_char(cp, '•')

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
