"""
Initial Setup: Create a presentation with a scatter chart on slide 4 (no trend line)
Task ID: impress_tct_057
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ──
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Revenue Trend Analysis"
    slide1.placeholders[1].text = "Finance Division — Q1 2025 Performance Review"

    # ── Slide 2: Executive Summary ──
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    tf2 = slide2.placeholders[1].text_frame
    tf2.clear()
    bullet_items = [
        "Total revenue grew 18% year-over-year in Q1 2025",
        "Marketing spend showed strong positive correlation with new customer acquisition",
        "Customer acquisition cost decreased from $142 to $118 per customer",
        "Regression analysis on next slide shows projected growth trajectory",
        "Recommendation: increase marketing budget allocation by 12% for Q2",
    ]
    for i, item in enumerate(bullet_items):
        if i == 0:
            tf2.paragraphs[0].text = item
        else:
            p = tf2.add_paragraph()
            p.text = item
            p.level = 0

    # ── Slide 3: Raw Data Table ──
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    txTitle.text_frame.paragraphs[0].text = "Marketing Spend vs. New Customers Acquired"
    run = txTitle.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Data for scatter chart (marketing spend in $K vs new customers)
    data_points = [
        (12.5, 48),
        (18.0, 67),
        (22.3, 82),
        (27.1, 95),
        (31.8, 110),
        (35.4, 128),
        (40.2, 139),
        (44.7, 155),
        (49.0, 168),
        (55.3, 192),
    ]

    headers = ["Month", "Marketing Spend ($K)", "New Customers"]
    months = ["Apr 2024", "May 2024", "Jun 2024", "Jul 2024", "Aug 2024",
              "Sep 2024", "Oct 2024", "Nov 2024", "Dec 2024", "Jan 2025"]

    rows = len(data_points) + 1
    cols = 3
    tbl_shape = slide3.shapes.add_table(rows, cols, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
    tbl = tbl_shape.table

    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1F497D'})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)

    for r, (spend, custs) in enumerate(data_points):
        tbl.cell(r + 1, 0).text = months[r]
        tbl.cell(r + 1, 1).text = f"${spend:.1f}K"
        tbl.cell(r + 1, 2).text = str(custs)

    # ── Slide 4: Scatter Chart ──
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txTitle4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.7))
    txTitle4.text_frame.paragraphs[0].text = "Marketing Spend vs. Customer Acquisition"
    run4 = txTitle4.text_frame.paragraphs[0].runs[0]
    run4.font.size = Pt(26)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    chart_data = XyChartData()
    series = chart_data.add_series('New Customers')
    for spend, custs in data_points:
        series.add_data_point(spend, custs)

    chart_frame = slide4.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(0.8), Inches(1.0),
        Inches(10.5), Inches(5.8),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Style the series
    plot = chart.plots[0]
    ser = plot.series[0]
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)

    # Axis labels via XML
    from lxml import etree
    chart_part = chart.part
    chart_xml = chart_part._element

    ns = {'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
          'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

    # Add X axis title
    catAx = chart_xml.find('.//c:valAx', ns)
    if catAx is not None:
        # There should be two valAx for scatter - first is X (catAx), second is Y (valAx)
        val_axes = chart_xml.findall('.//c:valAx', ns)
        if len(val_axes) >= 1:
            x_ax = val_axes[0]
            title_xml = etree.SubElement(x_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}title')
            tx = etree.SubElement(title_xml, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
            rich = etree.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
            bodyPr = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            lstStyle = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            p = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            t = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            t.text = "Marketing Spend ($K)"

        if len(val_axes) >= 2:
            y_ax = val_axes[1]
            title_xml = etree.SubElement(y_ax, '{http://schemas.openxmlformats.org/drawingml/2006/chart}title')
            tx = etree.SubElement(title_xml, '{http://schemas.openxmlformats.org/drawingml/2006/chart}tx')
            rich = etree.SubElement(tx, '{http://schemas.openxmlformats.org/drawingml/2006/chart}rich')
            bodyPr = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            lstStyle = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}lstStyle')
            p = etree.SubElement(rich, '{http://schemas.openxmlformats.org/drawingml/2006/main}p')
            r = etree.SubElement(p, '{http://schemas.openxmlformats.org/drawingml/2006/main}r')
            t = etree.SubElement(r, '{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            t.text = "New Customers"

    # ── Slide 5: Conclusions ──
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Takeaways & Next Steps"
    tf5 = slide5.placeholders[1].text_frame
    tf5.clear()
    conclusions = [
        "Strong linear relationship between marketing investment and customer growth",
        "Current trajectory suggests ~200 new customers at $60K monthly spend",
        "Diminishing returns may emerge beyond $70K — monitor closely",
        "Propose A/B testing of channel mix to optimize cost per acquisition",
        "Next review scheduled for end of Q2 2025",
    ]
    for i, item in enumerate(conclusions):
        if i == 0:
            tf5.paragraphs[0].text = item
        else:
            p = tf5.add_paragraph()
            p.text = item
            p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
