"""
Initial Setup: Fix misaligned rectangles on slide 5
Task ID: impress_fix_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Process Flow Overview"
    slide1.placeholders[1].text = "Q1 2025 Operations Review\nPrepared by Lisa Park"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Project Background"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "The operations team has identified five key process stages that drive our quarterly output."
    p2 = body2.add_paragraph()
    p2.text = "Each stage has been mapped to a responsible department with clear KPIs."
    p2.level = 1
    p3 = body2.add_paragraph()
    p3.text = "This presentation outlines the current flow and highlights areas for improvement."
    p3.level = 1

    # --- Slide 3: Key Metrics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Metrics Summary"
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Add a small table with metrics
    tbl_shape = slide3.shapes.add_table(4, 3, Inches(1), Inches(1.5), Inches(8), Inches(2.5))
    tbl = tbl_shape.table
    headers = ["Stage", "Avg Duration (days)", "Completion Rate"]
    data = [
        ["Intake & Triage", "2.3", "97%"],
        ["Analysis & Planning", "5.1", "89%"],
        ["Execution & Delivery", "8.7", "82%"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            tbl.cell(r, c).text = val

    # --- Slide 4: Team Responsibilities ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Team Responsibilities"
    body4 = slide4.placeholders[1].text_frame
    items = [
        "Intake: Customer Success (Sarah Chen, Marcus Davis)",
        "Analysis: Business Intelligence (Priya Sharma)",
        "Planning: Project Management (Tom Garcia)",
        "Execution: Engineering (Wei Zhang, Aisha Okafor)",
        "Review: Quality Assurance (James Wilson)",
    ]
    body4.text = items[0]
    for item in items[1:]:
        p = body4.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 5: Process Flow Diagram (misaligned rectangles) ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Title for slide 5
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(0.6))
    tf5 = title_box.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = "Process Flow Stages"
    r5 = p5.runs[0]
    r5.font.size = Pt(24)
    r5.font.bold = True
    r5.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # 5 rectangles representing process stages - MISALIGNED on purpose
    rect_width = Inches(2.0)
    rect_height = Inches(1.2)

    # These are deliberately uneven/random horizontal positions and slightly varied vertical
    labels = ["Intake", "Analysis", "Planning", "Execution", "Review"]
    colors = [
        RGBColor(0x4A, 0x90, 0xD9),  # blue
        RGBColor(0x5C, 0xB8, 0x5C),  # green
        RGBColor(0xF0, 0xAD, 0x4E),  # orange
        RGBColor(0xD9, 0x53, 0x4F),  # red
        RGBColor(0x9B, 0x59, 0xB6),  # purple
    ]

    # Misaligned positions - uneven horizontal spacing and slightly different vertical positions
    positions = [
        (Inches(0.3), Inches(3.2)),    # too far left, slightly high
        (Inches(3.1), Inches(3.5)),    # gap too large
        (Inches(5.0), Inches(2.9)),    # too close to previous, too high
        (Inches(8.5), Inches(3.6)),    # big gap
        (Inches(10.2), Inches(3.0)),   # close to right edge, high
    ]

    for i, (left, top) in enumerate(positions):
        shape = slide5.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, rect_width, rect_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = labels[i]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 6: Next Steps ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Next Steps"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Standardize the intake form by April 30, 2025"
    steps = [
        "Automate the analysis stage with BI dashboards",
        "Reduce execution cycle time by 15% in Q2",
        "Implement peer review checklist for QA stage",
        "Monthly process review meetings starting May 2025",
    ]
    for s in steps:
        p = body6.add_paragraph()
        p.text = s
        p.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
