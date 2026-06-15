"""
Reward Script: Align 5 rectangles on slide 5 evenly distributed horizontally
Task ID: impress_fix_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.4): All 5 rectangles share the same vertical position (top)
  Component 2 (0.3): Rectangles are evenly distributed horizontally (equal gaps)
  Component 3 (0.3): Shapes span the slide symmetrically with good coverage
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_067'

# Expected rectangle labels in order left-to-right
EXPECTED_LABELS = ['Intake', 'Analysis', 'Planning', 'Execution', 'Review']
NUM_RECTS = 5

# Tolerance for position comparisons (0.5% relative)
RELATIVE_TOL = 0.005
# Absolute tolerance for small values (in EMU, ~0.01 inch)
ABS_TOL = 9144


def is_approx_equal(val1, val2, rel_tol=RELATIVE_TOL, abs_tol=ABS_TOL):
    """Check if two EMU values are approximately equal."""
    if val1 == val2:
        return True
    diff = abs(val1 - val2)
    if diff <= abs_tol:
        return True
    max_val = max(abs(val1), abs(val2))
    if max_val == 0:
        return diff <= abs_tol
    return diff / max_val <= rel_tol


def get_rectangles(slide):
    """Extract AUTO_SHAPE rectangles from slide, sorted by left position."""
    rects = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            text = shape.text.strip() if hasattr(shape, 'text') and shape.text else ''
            rects.append({
                'text': text,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
            })
    # Sort by left position
    rects.sort(key=lambda r: r['left'])
    return rects


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 5 slides
    if len(prs.slides) < 5:
        print(f"CRITICAL: Presentation has {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[4]  # slide 5 (0-indexed)
    rects = get_rectangles(slide)

    # Precondition: exactly 5 rectangles on slide 5
    if len(rects) != NUM_RECTS:
        print(f"CRITICAL: Found {len(rects)} rectangles on slide 5, expected {NUM_RECTS}")
        print("REWARD: 0.0")
        return 0.0

    print(f"Found {len(rects)} rectangles on slide 5:")
    for i, r in enumerate(rects):
        print(f"  Rect {i}: text='{r['text']}', left={r['left']}, top={r['top']}, "
              f"width={r['width']}, height={r['height']}")

    # Component 1: All 5 rectangles have the same top position (0.4 points)
    # This verifies vertical alignment - all shapes at same Y coordinate
    try:
        tops = [r['top'] for r in rects]
        ref_top = tops[0]
        all_same_top = all(is_approx_equal(t, ref_top) for t in tops)

        if all_same_top:
            print(f"PASS: Component 1 - All rectangles aligned vertically at top={ref_top} (0.4 pts)")
            total_score += 0.4
        else:
            top_range = max(tops) - min(tops)
            print(f"FAIL: Component 1 - Tops not aligned. Values: {tops}, range={top_range} EMU")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Rectangles are evenly distributed horizontally (0.3 points)
    # Equal gap between consecutive rectangles
    try:
        # Calculate gaps between consecutive rectangles
        gaps = []
        for i in range(len(rects) - 1):
            gap = rects[i + 1]['left'] - (rects[i]['left'] + rects[i]['width'])
            gaps.append(gap)

        if len(gaps) >= 1:
            avg_gap = sum(gaps) / len(gaps)
            all_equal_gaps = all(is_approx_equal(g, avg_gap) for g in gaps)

            if all_equal_gaps and avg_gap > 0:
                print(f"PASS: Component 2 - Equal horizontal spacing, gap={avg_gap:.0f} EMU (0.3 pts)")
                total_score += 0.3
            else:
                print(f"FAIL: Component 2 - Unequal gaps: {[f'{g:.0f}' for g in gaps]}, avg={avg_gap:.0f}")
        else:
            print(f"FAIL: Component 2 - Could not compute gaps")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Shapes span the full slide width symmetrically (0.3 points)
    # The evenly distributed shapes should be positioned symmetrically relative to slide center
    # This checks that the distribution covers the slide properly, not bunched to one side
    try:
        slide_width = prs.slide_width
        # Check: leftmost shape's left edge and rightmost shape's right edge are
        # approximately symmetric around the slide center
        leftmost_left = rects[0]['left']
        rightmost_right = rects[-1]['left'] + rects[-1]['width']
        left_margin = leftmost_left
        right_margin = slide_width - rightmost_right

        # Margins should be approximately equal (symmetric distribution)
        margins_symmetric = is_approx_equal(left_margin, right_margin, rel_tol=0.05)

        # Also verify the shapes actually span a reasonable portion of the slide
        total_span = rightmost_right - leftmost_left
        coverage_ratio = total_span / slide_width
        good_coverage = coverage_ratio > 0.7  # shapes should span at least 70% of slide width

        if margins_symmetric and good_coverage:
            print(f"PASS: Component 3 - Symmetric distribution. Left margin={left_margin}, "
                  f"right margin={right_margin}, coverage={coverage_ratio:.2%} (0.3 pts)")
            total_score += 0.3
        else:
            if not margins_symmetric:
                print(f"FAIL: Component 3 - Asymmetric margins. Left={left_margin}, Right={right_margin}")
            if not good_coverage:
                print(f"FAIL: Component 3 - Low coverage: {coverage_ratio:.2%}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
