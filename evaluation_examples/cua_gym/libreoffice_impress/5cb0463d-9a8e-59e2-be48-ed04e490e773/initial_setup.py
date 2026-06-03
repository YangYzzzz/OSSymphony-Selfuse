"""
Initial Setup: Executive Dashboard presentation with 8 slides.
Slide 2 has title 'Key Metrics at a Glance' but no KPI content shapes.
Task ID: impress_exec_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title(slide, text):
    """Set title text on a slide if it has a title placeholder."""
    if slide.shapes.title:
        slide.shapes.title.text = text


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Executive Dashboard"
    slide1.placeholders[1].text = "Q4 2025 Performance Review\nPrepared by Strategic Planning Team"

    # --- Slide 2: Key Metrics at a Glance (EMPTY - no KPI shapes) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box
    add_textbox(slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Key Metrics at a Glance", font_size=32, bold=True,
                color=RGBColor(0x33, 0x33, 0x33), alignment=PP_ALIGN.LEFT)

    # --- Slide 3: Revenue Breakdown ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide3, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Revenue Breakdown by Region", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    # Revenue data table
    table_shape = slide3.shapes.add_table(6, 4, Inches(0.8), Inches(1.5), Inches(11), Inches(3.5))
    table = table_shape.table
    headers = ["Region", "Q3 Revenue", "Q4 Revenue", "Growth %"]
    data = [
        ["North America", "$28.4M", "$31.2M", "+9.9%"],
        ["Europe", "$14.7M", "$16.1M", "+9.5%"],
        ["Asia Pacific", "$8.3M", "$10.8M", "+30.1%"],
        ["Latin America", "$3.2M", "$4.4M", "+37.5%"],
        ["Total", "$54.6M", "$62.5M", "+14.5%"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
    for r, row in enumerate(data, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val

    # --- Slide 4: Customer Growth ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide4, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Customer Growth Trajectory", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_textbox(slide4, Inches(0.8), Inches(1.5), Inches(10), Inches(1.2),
                "Total active customers reached 1,247 in Q4 2025, representing a 18.3% increase "
                "year-over-year. Enterprise segment grew by 23%, while SMB added 156 new accounts.",
                font_size=16)
    # Customer milestone data
    milestones = [
        ("Q1 2025", "987 customers"),
        ("Q2 2025", "1,054 customers"),
        ("Q3 2025", "1,138 customers"),
        ("Q4 2025", "1,247 customers"),
    ]
    y_pos = Inches(3.2)
    for label, val in milestones:
        add_textbox(slide4, Inches(1.5), y_pos, Inches(3), Inches(0.5), label, font_size=14, bold=True)
        add_textbox(slide4, Inches(5.0), y_pos, Inches(3), Inches(0.5), val, font_size=14)
        y_pos += Inches(0.6)

    # --- Slide 5: Net Promoter Score ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide5, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Net Promoter Score Analysis", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_textbox(slide5, Inches(0.8), Inches(1.5), Inches(10), Inches(1.5),
                "Our NPS improved to 81 this quarter, placing us in the 'Excellent' category. "
                "Key drivers include improved onboarding experience (+12 points), enhanced "
                "customer support responsiveness (+8 points), and new product features (+5 points).",
                font_size=16)
    # NPS breakdown
    nps_data = [
        ("Promoters (9-10)", "68%"),
        ("Passives (7-8)", "19%"),
        ("Detractors (0-6)", "13%"),
    ]
    y_pos = Inches(3.5)
    for label, val in nps_data:
        add_textbox(slide5, Inches(1.5), y_pos, Inches(4), Inches(0.5), label, font_size=14)
        add_textbox(slide5, Inches(6.0), y_pos, Inches(2), Inches(0.5), val, font_size=14, bold=True)
        y_pos += Inches(0.6)

    # --- Slide 6: ARR Growth ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Annual Recurring Revenue Growth", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    add_textbox(slide6, Inches(0.8), Inches(1.5), Inches(10), Inches(1.5),
                "ARR growth accelerated to 24% year-over-year, driven by expansion revenue "
                "from existing customers and improved sales efficiency. Net revenue retention "
                "reached 118%, with gross churn declining to 4.2%.",
                font_size=16)
    arr_data = [
        ("Beginning ARR", "$50.4M"),
        ("New Business", "$8.1M"),
        ("Expansion", "$6.8M"),
        ("Churn", "-$2.8M"),
        ("Ending ARR", "$62.5M"),
    ]
    y_pos = Inches(3.5)
    for label, val in arr_data:
        add_textbox(slide6, Inches(1.5), y_pos, Inches(4), Inches(0.5), label, font_size=14)
        add_textbox(slide6, Inches(6.0), y_pos, Inches(2), Inches(0.5), val, font_size=14, bold=True)
        y_pos += Inches(0.6)

    # --- Slide 7: Strategic Initiatives ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide7, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Strategic Initiatives Status", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    initiatives = [
        "Platform Migration to Cloud-Native Architecture - 78% Complete",
        "Enterprise Sales Team Expansion (12 new AEs hired) - On Track",
        "Customer Success Program Redesign - Launched Q3",
        "APAC Market Entry Strategy - Phase 2 Underway",
        "AI-Powered Analytics Feature Set - Beta Testing",
    ]
    y_pos = Inches(1.5)
    for item in initiatives:
        add_textbox(slide7, Inches(1.0), y_pos, Inches(10), Inches(0.5), f"• {item}", font_size=14)
        y_pos += Inches(0.7)

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    add_textbox(slide8, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                "Next Steps & Q1 2026 Priorities", font_size=28, bold=True,
                color=RGBColor(0x33, 0x33, 0x33))
    next_steps = [
        "Finalize Q1 hiring plan for Product and Engineering teams",
        "Launch self-serve onboarding for SMB customers",
        "Complete SOC 2 Type II certification",
        "Roll out AI analytics to all Enterprise customers",
        "Expand partner ecosystem with 5 new integrations",
    ]
    y_pos = Inches(1.5)
    for i, item in enumerate(next_steps, 1):
        add_textbox(slide8, Inches(1.0), y_pos, Inches(10), Inches(0.5), f"{i}. {item}", font_size=14)
        y_pos += Inches(0.7)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
