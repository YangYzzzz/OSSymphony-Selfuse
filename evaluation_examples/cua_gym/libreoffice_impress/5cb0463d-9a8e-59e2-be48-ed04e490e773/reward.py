"""
Reward Script: KPI dashboard with four rounded rectangle shapes on slide 2
Task ID: impress_exec_040
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Four rounded rectangle shapes exist on slide 2
  Component 2 (0.25): Each shape is 3in x 2in (within tolerance)
  Component 3 (0.25): Correct fill colors paired with correct metric labels
  Component 4 (0.25): Correct text formatting (labels 14pt, values 28pt bold) with correct values
"""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_exec_040'

# Expected KPI cards: label -> (value, fill_color_hex)
EXPECTED_CARDS = {
    'Revenue':    ('$62.5M', 'E3F2FD'),
    'Customers':  ('1,247',  'E8F5E9'),
    'NPS':        ('81',     'FFF3E0'),
    'ARR Growth': ('24%',    'FCE4EC'),
}

EXPECTED_WIDTH = Inches(3)   # 2743200 EMU
EXPECTED_HEIGHT = Inches(2)  # 1828800 EMU
SIZE_TOLERANCE = 0.05  # 5% relative tolerance


def is_approx(val, expected, tol=SIZE_TOLERANCE):
    """Check if val is approximately equal to expected within relative tolerance."""
    if expected == 0:
        return val == 0
    return abs(val - expected) / expected <= tol


def get_rounded_rectangles(slide):
    """Find all rounded rectangle auto-shapes on a slide."""
    rects = []
    for shape in slide.shapes:
        try:
            if shape.auto_shape_type is not None and shape.auto_shape_type == 5:  # ROUNDED_RECTANGLE
                rects.append(shape)
        except (ValueError, AttributeError):
            continue
    return rects


def get_shape_fill_color(shape):
    """Get the solid fill color hex string of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def get_shape_text_info(shape):
    """Extract label and value text with font properties from a shape's text frame.
    Returns: list of (text, size_pt, is_bold) for each non-empty run.
    """
    runs_info = []
    if not hasattr(shape, 'text_frame'):
        return runs_info
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            text = (run.text or '').strip()
            if not text:
                continue
            size_pt = run.font.size / 12700 if run.font.size else None
            bold = run.font.bold if run.font.bold is not None else False
            runs_info.append((text, size_pt, bold))
    return runs_info


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Verify we have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[1]  # slide 2 (0-indexed)
    rects = get_rounded_rectangles(slide)

    # Component 1: Four rounded rectangle shapes exist on slide 2 (0.25 points)
    try:
        num_rects = len(rects)
        if num_rects == 4:
            print(f"PASS: Component 1 -- Found 4 rounded rectangles on slide 2 (0.25 pts)")
            total_score += 0.25
        elif num_rects >= 1:
            partial = 0.25 * (min(num_rects, 4) / 4)
            print(f"PARTIAL: Component 1 -- Found {num_rects} rounded rectangles (expected 4), awarding {partial:.3f} pts")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Found {num_rects} rounded rectangles on slide 2 (expected 4)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    if num_rects == 0:
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {total_score}")
        return total_score

    # Component 2: Each shape is 3in x 2in (0.25 points)
    try:
        correct_size_count = 0
        for rect in rects:
            w_ok = is_approx(rect.width, EXPECTED_WIDTH)
            h_ok = is_approx(rect.height, EXPECTED_HEIGHT)
            if w_ok and h_ok:
                correct_size_count += 1
            else:
                print(f"  INFO: {rect.name} size={rect.width/914400:.2f}x{rect.height/914400:.2f}in (expected 3.00x2.00in)")

        if correct_size_count == len(rects) and correct_size_count >= 4:
            print(f"PASS: Component 2 -- All 4 rectangles are 3in x 2in (0.25 pts)")
            total_score += 0.25
        elif correct_size_count > 0:
            partial = 0.25 * (correct_size_count / 4)
            print(f"PARTIAL: Component 2 -- {correct_size_count}/{len(rects)} rectangles have correct size ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No rectangles have the correct 3in x 2in size")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct fill colors paired with correct metric labels (0.25 points)
    try:
        matched_cards = 0
        for rect in rects:
            fill_color = get_shape_fill_color(rect)
            runs = get_shape_text_info(rect)
            # The first run is the label
            label = runs[0][0] if len(runs) >= 1 else None

            if label and label in EXPECTED_CARDS:
                expected_color = EXPECTED_CARDS[label][1]
                if fill_color and fill_color.upper() == expected_color.upper():
                    matched_cards += 1
                    print(f"  INFO: '{label}' has correct fill color #{fill_color}")
                else:
                    print(f"  INFO: '{label}' fill color #{fill_color} != expected #{expected_color}")
            else:
                print(f"  INFO: Unrecognized label '{label}' with fill #{fill_color}")

        if matched_cards == 4:
            print(f"PASS: Component 3 -- All 4 KPI cards have correct label-color pairing (0.25 pts)")
            total_score += 0.25
        elif matched_cards > 0:
            partial = 0.25 * (matched_cards / 4)
            print(f"PARTIAL: Component 3 -- {matched_cards}/4 cards have correct label-color pairing ({partial:.3f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No cards have correct label-color pairing")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct text formatting and values (0.25 points)
    # Labels should be 14pt (not bold), values should be 28pt bold
    try:
        format_score = 0
        points_per_card = 0.25 / 4  # 0.0625 per card

        for rect in rects:
            runs = get_shape_text_info(rect)
            if len(runs) < 2:
                print(f"  INFO: {rect.name} has {len(runs)} text runs (expected >=2: label + value)")
                continue

            label_text, label_size, label_bold = runs[0]
            value_text, value_size, value_bold = runs[1]

            card_issues = 0
            card_details = []

            # Check label is in expected set
            if label_text not in EXPECTED_CARDS:
                card_issues += 1
                card_details.append(f"label '{label_text}' not recognized")
            else:
                expected_value = EXPECTED_CARDS[label_text][0]
                # Check value text matches
                if value_text != expected_value:
                    card_issues += 1
                    card_details.append(f"value '{value_text}' != expected '{expected_value}'")

            # Check label font size ~14pt
            if label_size is not None and abs(label_size - 14.0) <= 1.0:
                card_details.append(f"label size {label_size}pt OK")
            else:
                card_issues += 1
                card_details.append(f"label size {label_size}pt != 14pt")

            # Check value font size ~28pt
            if value_size is not None and abs(value_size - 28.0) <= 1.0:
                card_details.append(f"value size {value_size}pt OK")
            else:
                card_issues += 1
                card_details.append(f"value size {value_size}pt != 28pt")

            # Check value is bold
            if value_bold:
                card_details.append("value bold OK")
            else:
                card_issues += 1
                card_details.append("value not bold")

            if card_issues == 0:
                format_score += points_per_card
                print(f"  PASS: '{label_text}' formatting correct: {'; '.join(card_details)}")
            else:
                print(f"  FAIL: '{label_text}' formatting issues: {'; '.join(card_details)}")

        if format_score >= 0.25 - 0.001:
            print(f"PASS: Component 4 -- All cards have correct formatting and values (0.25 pts)")
            total_score += 0.25
        elif format_score > 0:
            print(f"PARTIAL: Component 4 -- Some cards have correct formatting ({format_score:.3f} pts)")
            total_score += format_score
        else:
            print(f"FAIL: Component 4 -- No cards have correct formatting")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
