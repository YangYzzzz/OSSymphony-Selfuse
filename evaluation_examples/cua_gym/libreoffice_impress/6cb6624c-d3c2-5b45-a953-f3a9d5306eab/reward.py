"""
FINAL REWARD SCRIPT - SUCCESS
Task: While cleaning up the deck, I spotted that slide 44 still has the wrong heading. Could you rename that slide so the title reads exactly "Methodology" in LibreOffice Impress?
Generated: 2025-09-10 12:23:22
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_slide_44_title(file_path: str) -> float:
    """Verify that slide 44 exists and its title is exactly 'Methodology'.

    Returns a progressive score between 0.0 and 1.0:
        0.5 points  – Slide 44 exists
        0.5 points  – Slide 44 title text matches exactly 'Methodology'
    The final reward is printed as `REWARD: X.X` and returned as a float.
    """

    print(f"Verifying presentation file: {file_path}\n")
    score = 0.0
    max_score = 1.0
    SLIDE_EXISTS_WT   = 0.5  # weight for slide existence
    TITLE_CORRECT_WT  = 0.5  # weight for correct title text
    expected_title = "Methodology"

    # ---------- Prerequisite: file must exist and be loadable ----------
    if not os.path.exists(file_path):
        print("✗ File not found – cannot verify task\n")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Error loading presentation: {e}\n")
        print("REWARD: 0.0")
        return 0.0

    total_slides = len(prs.slides)
    print(f"Total slides detected: {total_slides}\n")

    # ---------- Requirement 1: Slide 44 must exist ----------
    if total_slides >= 44:
        score += SLIDE_EXISTS_WT
        print(f"✓ Slide 44 exists ({SLIDE_EXISTS_WT} points)\n")
        slide_44 = prs.slides[43]  # zero-indexed access
    else:
        print("✗ Slide 44 does not exist (0 points)\n")
        print(f"REWARD: {score}")
        return score  # cannot continue title check if slide missing

    # ---------- Requirement 2: Title must be exactly 'Methodology' ----------
    title_matches = False
    for shape in slide_44.shapes:
        # Many shapes can contain text; we check them all
        if hasattr(shape, "text"):
            text = shape.text.replace("\xa0", " ").strip()
            if text:
                print(f"  Found text on slide 44: '{text}'")
            if text == expected_title:
                title_matches = True
                break

    if title_matches:
        score += TITLE_CORRECT_WT
        print(f"✓ Slide 44 title matches expected '{expected_title}' ({TITLE_CORRECT_WT} points)\n")
    else:
        print(f"✗ Slide 44 title does not match '{expected_title}' (0 points)\n")

    # ---------- Final scoring ----------
    final_score = min(score, max_score)  # safety cap at 1.0
    print(f"Total score: {final_score}\n")
    print(f"REWARD: {final_score}")
    return final_score


# ------------------------------
# Execute verification when run directly
# ------------------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/while_cleaning_up_the_deck_i_spotted_that_slide_44_still_has_the_wrong_heading_could_you_rename_that_golden.pptx"
    verify_slide_44_title(FILE_PATH)
