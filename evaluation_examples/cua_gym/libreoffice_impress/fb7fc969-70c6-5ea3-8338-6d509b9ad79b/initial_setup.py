"""
Initial Setup: Create a 9-slide competitive analysis presentation with a table on slide 6.
Task ID: impress_gf2_040
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_040'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=None, alignment=PP_ALIGN.LEFT):
    """Helper to add a text box with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def add_bullet_slide(slide, title_text, bullets, title_color=RGBColor(0x1B, 0x3A, 0x5C)):
    """Add title and bulleted content to a slide."""
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(1.0),
                title_text, font_size=28, bold=True, color=title_color)
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(8.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    LIGHT_GRAY_BG = RGBColor(0xF2, 0xF2, 0xF2)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_textbox(slide1, Inches(1.0), Inches(2.0), Inches(8.0), Inches(1.5),
                "Q2 2025 Competitive Analysis Report", font_size=36, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.0), Inches(3.8), Inches(8.0), Inches(1.0),
                "Strategic Planning Division — CloudSync Technologies",
                font_size=18, color=RGBColor(0xB0, 0xC4, 0xDE), alignment=PP_ALIGN.CENTER)
    add_textbox(slide1, Inches(1.0), Inches(5.0), Inches(8.0), Inches(0.6),
                "Prepared by: Elena Vasquez, VP of Strategy | June 2025",
                font_size=14, color=RGBColor(0x88, 0xA8, 0xC8), alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide2, "Executive Summary", [
        "Market share shifted 3.2% in our favor during Q1 2025",
        "Key competitor NexaFlow lost two enterprise accounts to CloudSync",
        "DataBridge launched aggressive pricing strategy cutting rates by 22%",
        "Our NPS score of 72 leads the segment by 15 points",
        "Recommendation: accelerate AI integration roadmap to maintain lead",
    ])

    # --- Slide 3: Market Landscape ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide3, "Market Landscape Overview", [
        "Total addressable market grew to $14.8B (up 18% YoY)",
        "Enterprise segment accounts for 62% of total revenue",
        "Mid-market segment showing fastest growth at 31% CAGR",
        "Three new entrants identified in the SMB space",
        "Regulatory changes in EU creating compliance-driven demand",
    ])

    # --- Slide 4: Our Positioning ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide4, "CloudSync Market Position", [
        "Current market share: 24.6% (up from 21.4% in Q1 2024)",
        "Revenue: $3.64B annualized run rate",
        "Customer count: 2,847 enterprise accounts",
        "Average contract value increased to $1.28M",
        "Retention rate: 94.3% across all segments",
        "Key differentiators: AI-powered automation, real-time sync, compliance toolkit",
    ])

    # --- Slide 5: Methodology ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide5, "Analysis Methodology", [
        "Data sources: Gartner Magic Quadrant, Forrester Wave, internal CRM analytics",
        "Timeframe: January 2025 — May 2025 data collection period",
        "Scoring model: weighted matrix across 8 capability dimensions",
        "Peer review panel: 5 industry analysts validated scoring",
        "Confidence level: 92% based on data completeness assessment",
    ])

    # --- Slide 6: Competitive Comparison Table (NO ANIMATIONS) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide6, Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.8),
                "Feature Comparison Matrix", font_size=28, bold=True, color=DARK_BLUE)

    # Table: 7 rows (header + 6 competitors), 5 columns
    rows, cols = 7, 5
    tbl_shape = slide6.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.0)
    )
    table = tbl_shape.table

    # Set column widths
    table.columns[0].width = Inches(1.8)   # Company
    table.columns[1].width = Inches(1.8)   # AI Features
    table.columns[2].width = Inches(1.8)   # Scalability
    table.columns[3].width = Inches(1.8)   # Pricing
    table.columns[4].width = Inches(1.8)   # Support

    # Header row
    headers = ["Company", "AI Features", "Scalability", "Pricing ($/user/mo)", "Support Rating"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.color.rgb = WHITE
        # Dark blue header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = DARK_BLUE

    # Data rows
    data = [
        ["CloudSync (Us)",    "Advanced",  "Enterprise+", "$42",   "9.2 / 10"],
        ["NexaFlow",          "Moderate",  "Enterprise",  "$38",   "7.8 / 10"],
        ["DataBridge",        "Basic",     "Mid-Market",  "$29",   "8.1 / 10"],
        ["SyncWave Pro",      "Advanced",  "Enterprise",  "$45",   "8.5 / 10"],
        ["FlowMatrix AI",    "Advanced",  "Enterprise+", "$51",   "7.3 / 10"],
        ["CoreLink Systems",  "Moderate",  "Mid-Market",  "$34",   "8.9 / 10"],
    ]

    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = DARK_GRAY
            # Alternating row colors
            if r % 2 == 1:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = LIGHT_GRAY_BG

    # --- Slide 7: SWOT Analysis ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide7, "SWOT Analysis — CloudSync", [
        "Strengths: Industry-leading AI, highest NPS, strong brand recognition",
        "Weaknesses: Premium pricing limits SMB penetration, slow mobile app updates",
        "Opportunities: EU compliance mandate, mid-market expansion, AI copilot demand",
        "Threats: DataBridge price war, FlowMatrix AI feature parity, talent competition",
    ])

    # --- Slide 8: Strategic Recommendations ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide8, "Strategic Recommendations", [
        "1. Launch mid-market tier at $28/user/mo by Q3 2025",
        "2. Accelerate AI copilot release to maintain differentiation",
        "3. Invest $12M in customer success to defend retention rate",
        "4. Pursue strategic partnership with Salesforce for CRM integration",
        "5. Monitor DataBridge pricing closely — prepare counter-offer playbook",
    ])

    # --- Slide 9: Next Steps & Timeline ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_bullet_slide(slide9, "Next Steps & Timeline", [
        "June 15: Board presentation of competitive findings",
        "July 1: Mid-market pricing proposal finalized",
        "July 15: AI copilot beta launch to 50 enterprise clients",
        "August 1: Customer success team expansion (15 new hires)",
        "September 1: Q3 competitive landscape refresh",
        "Quarterly review cadence established with strategy team",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
