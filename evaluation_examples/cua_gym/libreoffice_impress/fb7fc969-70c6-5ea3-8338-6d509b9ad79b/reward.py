"""
Reward Script: Verify 'Build In' animation sequence on slide 6 table
Task ID: impress_gf2_040
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Slide 6 has a timing element with animation entries (none in initial)
  Component 2 (0.25): Exactly 7 entrance animations targeting the table (one per row)
  Component 3 (0.25): All animations use Wipe from Left (presetID=22, presetSubtype=1, filter=wipe(left))
  Component 4 (0.25): First animation is clickEffect, rows 2-7 are afterEffect with 600ms delay
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_gf2_040'

# Namespaces used in PPTX animation XML
NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def parse_animations(pptx_path, slide_idx=5):
    """
    Parse slide 6 (0-indexed=5) animation XML.
    Returns a list of dicts describing each animation entry, or None if no timing.
    """
    slide_file = f'ppt/slides/slide{slide_idx + 1}.xml'
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        try:
            with zf.open(slide_file) as f:
                root = ET.parse(f).getroot()
        except KeyError:
            return None

    timing = root.find('.//p:timing', NS)
    if timing is None:
        return None

    # Find all cTn elements with presetClass="entr" (entrance animations)
    animations = []
    for ctn in timing.iter():
        tag = ctn.tag.split('}')[-1] if '}' in ctn.tag else ctn.tag
        if tag != 'cTn':
            continue
        preset_class = ctn.get('presetClass')
        if preset_class != 'entr':
            continue

        preset_id = ctn.get('presetID')
        preset_subtype = ctn.get('presetSubtype')
        node_type = ctn.get('nodeType')

        # Find the target spid and paragraph range
        spid = None
        pRg_st = None
        pRg_end = None
        wipe_filter = None

        for sp_tgt in ctn.iter():
            sp_tag = sp_tgt.tag.split('}')[-1] if '}' in sp_tgt.tag else sp_tgt.tag
            if sp_tag == 'spTgt' and spid is None:
                spid = sp_tgt.get('spid')
            if sp_tag == 'pRg' and pRg_st is None:
                pRg_st = sp_tgt.get('st')
                pRg_end = sp_tgt.get('end')
            if sp_tag == 'animEffect' and wipe_filter is None:
                wipe_filter = sp_tgt.get('filter')

        # Find delay: look at parent cTn stCondLst for delay value
        # The delay is on the grandparent par's cTn element
        delay = None
        parent = ctn.find('.//p:stCondLst/p:cond', NS)
        # We need to look at the containing par's delay instead
        # The delay is in the wrapping cTn's stCondLst, not in this cTn
        # We'll extract it differently below

        animations.append({
            'preset_id': preset_id,
            'preset_subtype': preset_subtype,
            'preset_class': preset_class,
            'node_type': node_type,
            'spid': spid,
            'pRg_st': pRg_st,
            'pRg_end': pRg_end,
            'wipe_filter': wipe_filter,
        })

    # Now extract delays from the wrapping par elements
    # For afterEffect entries, the delay is on the sibling/parent cTn
    # Re-parse to get delays properly
    main_seq = timing.find('.//p:seq/p:cTn/p:childTnLst', NS)
    if main_seq is None:
        # Try without the namespace approach - iterate more carefully
        for elem in timing.iter():
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if tag == 'seq':
                for child in elem:
                    child_tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    if child_tag == 'cTn':
                        for sub in child:
                            sub_tag = sub.tag.split('}')[-1] if '}' in sub.tag else sub.tag
                            if sub_tag == 'childTnLst':
                                main_seq = sub
                                break
                break

    delays = []
    if main_seq is not None:
        for par_elem in main_seq:
            par_tag = par_elem.tag.split('}')[-1] if '}' in par_elem.tag else par_elem.tag
            if par_tag != 'par':
                continue
            # Each top-level par wraps one animation entry
            # Look for the inner delay on the middle cTn
            # Structure: par > cTn > childTnLst > par > cTn(delay) > childTnLst > par > cTn(animation)
            delay_val = None
            for inner in par_elem.iter():
                inner_tag = inner.tag.split('}')[-1] if '}' in inner.tag else inner.tag
                if inner_tag == 'cTn' and inner.get('fill') == 'hold' and inner.get('presetClass') is None:
                    # Look at stCondLst for delay
                    for cond in inner.iter():
                        cond_tag = cond.tag.split('}')[-1] if '}' in cond.tag else cond.tag
                        if cond_tag == 'cond' and cond.get('delay') is not None and cond.get('evt') is None:
                            d = cond.get('delay')
                            if d and d != '0':
                                delay_val = d
                                break
                    if delay_val:
                        break
            delays.append(delay_val)

    # Merge delays into animations
    for i, anim in enumerate(animations):
        if i < len(delays):
            anim['delay'] = delays[i]
        else:
            anim['delay'] = None

    return animations


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    if not os.path.exists(file_path):
        print(f"CRITICAL: File not found: {file_path}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file is a valid pptx with at least 6 slides
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        if len(prs.slides) < 6:
            print(f"CRITICAL: Only {len(prs.slides)} slides, need at least 6")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot load presentation: {e}")
        print("REWARD: 0.0")
        return 0.0

    animations = parse_animations(file_path, slide_idx=5)

    # Component 1: Slide 6 has entrance animations (0.25 points)
    # This fails on initial (no timing element) and passes on golden
    try:
        if animations is not None and len(animations) > 0:
            print(f"PASS: Component 1 — Slide 6 has {len(animations)} entrance animation(s) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — No entrance animations found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Exactly 7 entrance animations (one per table row) (0.25 points)
    # Award partial: 0.15 if 5-9 animations, full 0.25 if exactly 7
    try:
        if animations and len(animations) == 7:
            # Verify they target sequential paragraph ranges (0-6)
            pRg_values = [int(a['pRg_st']) for a in animations if a['pRg_st'] is not None]
            if sorted(pRg_values) == list(range(7)):
                print(f"PASS: Component 2 — Exactly 7 animations targeting paragraphs 0-6 (0.25 pts)")
                total_score += 0.25
            else:
                print(f"PARTIAL: Component 2 — 7 animations but paragraph targets are {pRg_values}, expected 0-6 (0.15 pts)")
                total_score += 0.15
        elif animations and 5 <= len(animations) <= 9:
            print(f"PARTIAL: Component 2 — Found {len(animations)} animations, expected 7 (0.1 pts)")
            total_score += 0.1
        else:
            count = len(animations) if animations else 0
            print(f"FAIL: Component 2 — Found {count} animations, expected 7")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: All animations use Wipe from Left (0.25 points)
    # presetID=22 (Wipe), presetSubtype=1 (from left), filter=wipe(left)
    try:
        if animations and len(animations) > 0:
            wipe_count = 0
            for anim in animations:
                is_wipe = (anim['preset_id'] == '22')
                is_from_left = (anim['preset_subtype'] == '1')
                has_filter = (anim['wipe_filter'] is not None and 'wipe' in str(anim['wipe_filter']).lower())
                if is_wipe and (is_from_left or has_filter):
                    wipe_count += 1

            ratio = wipe_count / len(animations)
            if ratio == 1.0:
                print(f"PASS: Component 3 — All {wipe_count}/{len(animations)} animations use Wipe from Left (0.25 pts)")
                total_score += 0.25
            elif ratio >= 0.5:
                pts = round(0.25 * ratio, 2)
                print(f"PARTIAL: Component 3 — {wipe_count}/{len(animations)} use Wipe from Left ({pts} pts)")
                total_score += pts
            else:
                print(f"FAIL: Component 3 — Only {wipe_count}/{len(animations)} use Wipe from Left")
        else:
            print(f"FAIL: Component 3 — No animations to check")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: First animation is clickEffect (On Click), rows 2-7 are afterEffect with 600ms delay (0.25 points)
    try:
        if animations and len(animations) >= 2:
            sub_score = 0.0

            # Check first animation is clickEffect
            first = animations[0]
            if first['node_type'] == 'clickEffect':
                print(f"  PASS: First animation is clickEffect (On Click)")
                sub_score += 0.1
            else:
                print(f"  FAIL: First animation nodeType={first['node_type']}, expected clickEffect")

            # Check remaining animations are afterEffect with 600ms delay
            after_count = 0
            delay_count = 0
            for anim in animations[1:]:
                if anim['node_type'] == 'afterEffect':
                    after_count += 1
                if anim.get('delay') == '600':
                    delay_count += 1

            expected_after = len(animations) - 1
            if expected_after > 0:
                if after_count == expected_after:
                    print(f"  PASS: All {after_count} subsequent animations are afterEffect")
                    sub_score += 0.075
                else:
                    print(f"  FAIL: {after_count}/{expected_after} subsequent animations are afterEffect")

                if delay_count == expected_after:
                    print(f"  PASS: All {delay_count} subsequent animations have 600ms delay")
                    sub_score += 0.075
                elif delay_count > 0:
                    ratio = delay_count / expected_after
                    pts = round(0.075 * ratio, 3)
                    print(f"  PARTIAL: {delay_count}/{expected_after} have 600ms delay ({pts} pts)")
                    sub_score += pts
                else:
                    print(f"  FAIL: No subsequent animations have 600ms delay")

            print(f"RESULT: Component 4 — trigger/delay checks ({round(sub_score, 3)} pts)")
            total_score += sub_score
        else:
            print(f"FAIL: Component 4 — Not enough animations to check triggers")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
