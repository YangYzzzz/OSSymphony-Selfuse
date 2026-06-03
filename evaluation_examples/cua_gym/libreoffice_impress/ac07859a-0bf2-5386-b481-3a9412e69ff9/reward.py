"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 3 there’s an image called “Picture 1” that I need resized so its height is exactly 20 cm (width can scale automatically). Separately, slide 6 should look uniform, so every text box on that slide has to use a 40 pt font. What steps do I follow in LibreOffice Impress to knock those two tweaks out?
Generated: 2025-09-10 14:23:54
Status: success
Model: azure-o3
Total Steps: 8
"""

import os
from pathlib import Path
from pptx import Presentation

# -----------------------------------------------------------------------------
# Reward verification script for LibreOffice Impress task
# -----------------------------------------------------------------------------
# Task requirements
#   1. On slide 3 an image named "Picture 1" must have its height set to EXACTLY
#      20 cm (width may scale automatically).
#   2. On slide 6 every text box must use a font size of 40 pt.
#
# The script awards 0.5 points for each correctly completed requirement, giving
# a progressive score from 0.0 to 1.0.  Only when BOTH conditions are fully met
# is the final reward exactly 1.0.
# -----------------------------------------------------------------------------

def emu_from_cm(cm: float) -> int:
    """Convert centimetres to English Metric Units (EMU)."""
    return int(cm * 360000)  # 1 cm  = 360 000 EMU

def emu_from_pt(pt: float) -> int:
    """Convert points to EMU."""
    return int(pt * 12700)   # 1 pt = 12 700 EMU

def verify_presentation(file_path: str) -> float:
    """Verify the presentation meets all task requirements and return a score."""

    max_score = 1.0
    score     = 0.0

    # ------------------------------------------------------------------
    # Load presentation file (no points for merely existing or loading)
    # ------------------------------------------------------------------
    try:
        if not Path(file_path).exists():
            print("✗ File does not exist:", file_path)
            return 0.0
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded successfully – {len(prs.slides)} slides found")
    except Exception as exc:
        print("✗ Error loading presentation:", exc)
        return 0.0

    # ------------------------------------------------------------------
    # Requirement 1 – Slide 3 image height = 20 cm
    # ------------------------------------------------------------------
    try:
        target_height_emu = emu_from_cm(20)
        tolerance_height = 1000  # EMU tolerance (≈0.003 cm)

        if len(prs.slides) < 3:
            print("✗ Slide 3 is missing (requirement 1)")
        else:
            slide3 = prs.slides[2]
            picture_found   = False
            correct_height  = False

            for shape in slide3.shapes:
                # PICTURE shape_type == 13 according to python-pptx constants
                if shape.shape_type == 13 and getattr(shape, "name", "").lower() == "picture 1":
                    picture_found = True
                    print(f"  • Found 'Picture 1' – height = {shape.height} EMU")
                    if abs(shape.height - target_height_emu) <= tolerance_height:
                        correct_height = True
                    break

            if picture_found and correct_height:
                print("✓ Requirement 1 fulfilled – image height is 20 cm (0.5 pts)")
                score += 0.5
            elif not picture_found:
                print("✗ Requirement 1 failed – 'Picture 1' not found on slide 3")
            else:
                print("✗ Requirement 1 failed – image height is incorrect")
    except Exception as exc:
        print("✗ Error verifying requirement 1:", exc)

    # ------------------------------------------------------------------
    # Requirement 2 – Slide 6 text boxes font size = 40 pt
    # ------------------------------------------------------------------
    try:
        target_font_emu = emu_from_pt(40)
        tolerance_font  = 100  # EMU tolerance (≈0.008 pt)

        if len(prs.slides) < 6:
            print("✗ Slide 6 is missing (requirement 2)")
        else:
            slide6       = prs.slides[5]
            total_runs   = 0
            correct_runs = 0

            for shape in slide6.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            total_runs += 1
                            size = run.font.size
                            # A run without explicit size fails the requirement
                            if size is not None and abs(size - target_font_emu) <= tolerance_font:
                                correct_runs += 1

            if total_runs == 0:
                print("✗ Requirement 2 failed – no text found on slide 6")
            else:
                print(f"  • Slide 6 font check: {correct_runs}/{total_runs} text runs at 40 pt")
                if correct_runs == total_runs:
                    print("✓ Requirement 2 fulfilled – all text is 40 pt (0.5 pts)")
                    score += 0.5
                else:
                    print("✗ Requirement 2 failed – not all text is 40 pt")
    except Exception as exc:
        print("✗ Error verifying requirement 2:", exc)

    # ------------------------------------------------------------------
    # Final score
    # ------------------------------------------------------------------
    final_score = min(score, max_score)
    print(f"Final score: {final_score}")
    return final_score

# -----------------------------------------------------------------------------
# Execute verification when script is run directly (this is required so the VM
# test harness executes the checks automatically).
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    TEST_FILE = "/home/user/on_slide_3_theres_an_image_called_picture_1_that_i_need_resized_so_its_height_is_exactly_20_cm_width_golden.pptx"
    reward = verify_presentation(TEST_FILE)
    print("REWARD:", reward)
