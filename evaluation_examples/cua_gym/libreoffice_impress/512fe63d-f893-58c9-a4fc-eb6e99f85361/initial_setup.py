"""
Initial Setup: Create a Price List presentation with tab-separated text on slide 2
Task ID: impress_gf1_042
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_042'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Cm(25.4)   # Standard 10 inches
    prs.slide_height = Cm(19.05)  # Standard 7.5 inches

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Gourmet Kitchen Supplies"
    slide1.placeholders[1].text = "Product Price List - Spring 2025"

    # --- Slide 2: Tab-separated price list (the target slide) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add a title text box at top
    title_box = slide2.shapes.add_textbox(Cm(1), Cm(0.5), Cm(23), Cm(2))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = "Product Price List"
    p_title.alignment = PP_ALIGN.CENTER
    run_title = p_title.runs[0]
    run_title.font.size = Pt(28)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Add the main text box with tab-separated data
    # Position it to span most of the slide width
    txBox = slide2.shapes.add_textbox(Cm(1), Cm(3), Cm(23.5), Cm(14))
    tf = txBox.text_frame
    tf.word_wrap = True

    # Header line
    header_data = "Item\tDescription\tPrice"
    # Data lines - realistic product data
    data_lines = [
        "Chef's Knife\tProfessional 8-inch stainless steel blade\t$89.95",
        "Copper Saucepan\tHand-hammered 2-quart with brass handle\t$145.00",
        "Bamboo Cutting Board\tExtra-large 18x24 inch end-grain\t$67.50",
        "Cast Iron Skillet\tPre-seasoned 12-inch Lodge Heritage\t$54.99",
        "Silicone Spatula Set\t5-piece heat-resistant up to 600F\t$24.95",
        "Digital Kitchen Scale\tPrecision 0.1g capacity 11lb max\t$32.00",
        "Stainless Mixing Bowls\tNesting set of 6 with non-slip base\t$49.95",
    ]

    all_lines = [header_data] + data_lines

    # First paragraph (header)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = all_lines[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Data paragraphs
    for line in all_lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # NOTE: No custom tab stops set - using default 1.25cm spacing
    # This causes the columns to be misaligned

    # --- Slide 3: Category Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    tb3 = slide3.shapes.add_textbox(Cm(2), Cm(1), Cm(21), Cm(2))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Product Categories"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.size = Pt(24)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    cat_box = slide3.shapes.add_textbox(Cm(2), Cm(4), Cm(21), Cm(12))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    categories = [
        ("Cutlery & Knives", "Premium blades from German and Japanese makers"),
        ("Cookware", "Cast iron, copper, and stainless steel essentials"),
        ("Prep & Storage", "Cutting boards, mixing bowls, and containers"),
        ("Gadgets & Tools", "Scales, thermometers, and specialty items"),
        ("Bakeware", "Sheet pans, cake molds, and silicone mats"),
    ]
    for i, (cat, desc) in enumerate(categories):
        if i == 0:
            p = tf_cat.paragraphs[0]
        else:
            p = tf_cat.add_paragraph()
        run = p.add_run()
        run.text = cat
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)

        p2 = tf_cat.add_paragraph()
        run2 = p2.add_run()
        run2.text = desc
        run2.font.size = Pt(12)
        run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        # Add spacing between categories
        if i < len(categories) - 1:
            spacer = tf_cat.add_paragraph()
            spacer.add_run().text = ""

    # --- Slide 4: Contact & Ordering ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    tb4 = slide4.shapes.add_textbox(Cm(2), Cm(1), Cm(21), Cm(2))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "How to Order"
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.runs[0]
    r4.font.size = Pt(24)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    info_box = slide4.shapes.add_textbox(Cm(2), Cm(4), Cm(21), Cm(10))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    info_lines = [
        "Online: www.gourmetkitchensupplies.com",
        "Phone: +1 (555) 234-5678",
        "Email: orders@gourmetkitchen.com",
        "",
        "Business Hours: Mon-Fri 8:00 AM - 6:00 PM EST",
        "Free shipping on orders over $100",
        "30-day satisfaction guarantee on all products",
    ]
    for i, line in enumerate(info_lines):
        if i == 0:
            p = tf_info.paragraphs[0]
        else:
            p = tf_info.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
