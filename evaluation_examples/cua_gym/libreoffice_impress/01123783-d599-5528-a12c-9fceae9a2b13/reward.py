"""
Reward Script: Align and distribute shapes on slide 4
Task ID: impress_ndo_046
Domain: libreoffice_impress
Scoring:
  Component 1 (0.5): All 3 shapes horizontally centered (center_x == slide_width/2)
  Component 2 (0.3): Equal vertical spacing between shapes
  Component 3 (0.2): Top-to-bottom order preserved (Rectangle, Circle, Triangle)
"""

import os
from pptx import Presentation
from pptx.util import Emu

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_046'


def persist_app_state(domain):
    """Save any unsaved GUI edits before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # 0-indexed, slide 4
    slide_width = prs.slide_width

    # Find the three target shapes by name
    shape_map = {}
    for shape in slide.shapes:
        if shape.name in ('Rectangle', 'Circle', 'Triangle'):
            center_x = shape.left + shape.width // 2
            center_y = shape.top + shape.height // 2
            shape_map[shape.name] = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'center_x': center_x,
                'center_y': center_y,
            }

    # Precondition: all three shapes must exist
    required = ['Rectangle', 'Circle', 'Triangle']
    for name in required:
        if name not in shape_map:
            print(f"FAIL: Shape '{name}' not found on slide 4")
            print("REWARD: 0.0")
            return 0.0

    expected_center_x = slide_width // 2  # 4572000 EMU for standard 10-inch slide

    # Component 1: All 3 shapes horizontally centered (0.5 points)
    # Each shape's center_x should equal slide_width / 2
    # Tolerance: 0.5% of slide width (~45720 EMU ~ 0.13cm)
    try:
        tolerance = slide_width * 0.005
        centered_count = 0
        for name in required:
            cx = shape_map[name]['center_x']
            diff = abs(cx - expected_center_x)
            if diff <= tolerance:
                centered_count += 1
                print(f"  PASS: {name} center_x={cx} (expected ~{expected_center_x}, diff={diff})")
            else:
                print(f"  FAIL: {name} center_x={cx} (expected ~{expected_center_x}, diff={diff})")

        if centered_count == 3:
            print(f"PASS: Component 1 -- All 3 shapes horizontally centered (0.5 pts)")
            total_score += 0.5
        elif centered_count >= 1:
            partial = round(0.5 * centered_count / 3, 2)
            print(f"PARTIAL: Component 1 -- {centered_count}/3 shapes centered ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- No shapes are horizontally centered")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Equal vertical spacing between shapes (0.3 points)
    # Sort shapes by center_y, then check that gaps between consecutive shapes are equal
    try:
        sorted_shapes = sorted(required, key=lambda n: shape_map[n]['center_y'])
        cy_vals = [shape_map[n]['center_y'] for n in sorted_shapes]

        gap1 = cy_vals[1] - cy_vals[0]
        gap2 = cy_vals[2] - cy_vals[1]

        # Tolerance: 2% of slide height for spacing equality
        spacing_tolerance = prs.slide_height * 0.02
        spacing_diff = abs(gap1 - gap2)

        print(f"  Vertical centers (sorted): {sorted_shapes[0]}={cy_vals[0]}, {sorted_shapes[1]}={cy_vals[1]}, {sorted_shapes[2]}={cy_vals[2]}")
        print(f"  Gap1={gap1}, Gap2={gap2}, diff={spacing_diff}, tolerance={spacing_tolerance}")

        if spacing_diff <= spacing_tolerance and gap1 > 0 and gap2 > 0:
            print(f"PASS: Component 2 -- Equal vertical spacing (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 -- Unequal vertical spacing (gap diff={spacing_diff})")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Top-to-bottom order preserved AND shapes are centered (0.2 points)
    # Compound check: order is only meaningful when centering has been applied.
    # This ensures the check FAILS on initial_env (where centering hasn't happened).
    try:
        rect_cy = shape_map['Rectangle']['center_y']
        circle_cy = shape_map['Circle']['center_y']
        tri_cy = shape_map['Triangle']['center_y']

        order_ok = rect_cy < circle_cy < tri_cy
        # Require at least 2 shapes centered to confirm alignment was attempted
        at_least_two_centered = centered_count >= 2

        if order_ok and at_least_two_centered:
            print(f"PASS: Component 3 -- Order preserved + shapes centered: Rectangle({rect_cy}) < Circle({circle_cy}) < Triangle({tri_cy}) (0.2 pts)")
            total_score += 0.2
        elif not order_ok:
            print(f"FAIL: Component 3 -- Order not preserved: Rectangle({rect_cy}), Circle({circle_cy}), Triangle({tri_cy})")
        else:
            print(f"FAIL: Component 3 -- Order preserved but shapes not centered ({centered_count}/3 centered)")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist app state before verification
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
