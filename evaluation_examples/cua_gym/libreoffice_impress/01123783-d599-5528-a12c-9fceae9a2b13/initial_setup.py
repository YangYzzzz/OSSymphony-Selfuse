"""
Initial Setup: Create presentation with three misaligned shapes on slide 4
Task ID: impress_ndo_046
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_046'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    # Standard widescreen: 33.867cm x 19.05cm (13.333 x 7.5 inches)
    slide_w = prs.slide_width   # 12192000 EMU
    slide_h = prs.slide_height  # 6858000 EMU

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Nexus Analytics - Q3 Strategy Review"
    if 1 in slide1.placeholders:
        slide1.placeholders[1].text = "Prepared by Rachel Morrison, VP of Operations"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Revenue growth of 18.3% year-over-year driven by enterprise segment"
    p = body2.add_paragraph()
    p.text = "Customer retention rate improved to 94.7% from 91.2%"
    p = body2.add_paragraph()
    p.text = "New product launches accounted for $12.4M in incremental revenue"
    p = body2.add_paragraph()
    p.text = "Operational efficiency gains reduced overhead by 7.8%"

    # --- Slide 3: Market Analysis ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Market Landscape & Competitive Position"
    body3 = slide3.placeholders[1].text_frame
    body3.text = "Total addressable market grew to $48.6B (up 12% from Q2)"
    p = body3.add_paragraph()
    p.text = "Nexus market share: 8.3% (ranked #3 in enterprise solutions)"
    p = body3.add_paragraph()
    p.text = "Key competitor Meridian Corp acquired DataFlow Inc for $2.1B"
    p = body3.add_paragraph()
    p.text = "Emerging opportunity in AI-driven analytics projected at $6.2B by 2027"

    # --- Slide 4: Layout slide with three randomly-placed shapes ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
    slide4.shapes.title.text = "Visual Layout Draft"

    # Shape dimensions
    rect_w, rect_h = Cm(4), Cm(3)
    circle_w, circle_h = Cm(3), Cm(3)
    tri_w, tri_h = Cm(3.5), Cm(3)

    # Rectangle at (3cm, 2cm) - randomly placed, NOT centered
    rect = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(3), Cm(2), rect_w, rect_h
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x41, 0x72, 0xC4)  # Blue
    rect.name = "Rectangle"
    tf = rect.text_frame
    tf.paragraphs[0].text = "Strategy"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # Circle at (15cm, 10cm) - randomly placed, NOT centered
    circle = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(15), Cm(10), circle_w, circle_h
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)  # Orange
    circle.name = "Circle"
    tf = circle.text_frame
    tf.paragraphs[0].text = "Growth"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(14)
        run.font.bold = True

    # Triangle at (8cm, 16cm) - randomly placed, NOT centered
    triangle = slide4.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE, Cm(8), Cm(16), tri_w, tri_h
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = RGBColor(0x70, 0xAD, 0x47)  # Green
    triangle.name = "Triangle"
    tf = triangle.text_frame
    tf.paragraphs[0].text = "Innovation"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(12)
        run.font.bold = True

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps & Action Items"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Finalize Q4 budget allocation by October 15th"
    p = body5.add_paragraph()
    p.text = "Launch enterprise pilot program with 5 key accounts"
    p = body5.add_paragraph()
    p.text = "Complete competitive analysis report for board review"
    p = body5.add_paragraph()
    p.text = "Schedule regional team alignment sessions"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
