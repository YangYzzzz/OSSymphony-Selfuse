"""
Initial Setup: Add animations to slide 3 bullet items and summary box
Task ID: impress_gf3_013
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_gf3_013'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Professional Development Workshop"
    slide1.placeholders[1].text = "Building Skills for the Modern Workplace\nQ2 2025 Series"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Workshop Agenda"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Introduction and Overview"
    for item in ["Current Industry Trends", "Best Practices Discussion",
                 "Hands-on Exercises", "Team Collaboration Session",
                 "Q&A and Wrap-up"]:
        p = body2.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 3: Best Practices (KEY SLIDE) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout for full control

    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Best Practices"
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.runs[0]
    run_title.font.size = Pt(36)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Bullet list with 5 items
    bullet_box = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(10.5), Inches(3.5))
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True

    bullet_items = [
        "Establish clear communication channels across all project teams",
        "Document decisions and rationale within 24 hours of key meetings",
        "Conduct weekly retrospectives to identify process improvements",
        "Use version control for all deliverables, not just source code",
        "Schedule dedicated focus time blocks to minimize context switching",
    ]

    for i, item in enumerate(bullet_items):
        if i == 0:
            p = tf_bullets.paragraphs[0]
        else:
            p = tf_bullets.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Summary box (rectangle) at the bottom
    summary_shape = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.0), Inches(5.5), Inches(7.0), Inches(1.2)
    )
    summary_shape.fill.solid()
    summary_shape.fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    summary_shape.line.color.rgb = RGBColor(0x14, 0x28, 0x44)

    tf_summary = summary_shape.text_frame
    tf_summary.word_wrap = True
    p_summary = tf_summary.paragraphs[0]
    p_summary.text = "Apply these daily"
    p_summary.alignment = PP_ALIGN.CENTER
    run_summary = p_summary.runs[0]
    run_summary.font.size = Pt(22)
    run_summary.font.bold = True
    run_summary.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 4: Case Study ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Case Study: Project Aurora"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "Timeline: January - March 2025"
    for item in ["Team of 12 cross-functional members",
                 "Delivered 3 weeks ahead of schedule",
                 "Client satisfaction score: 4.8/5.0",
                 "Key success factor: Consistent daily standups"]:
        p = body4.add_paragraph()
        p.text = item

    # --- Slide 5: Metrics ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Performance Metrics"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "On-time delivery rate: 94%"
    for item in ["Average sprint velocity increase: 23%",
                 "Defect rate reduction: 41% year-over-year",
                 "Team engagement score: 8.7/10",
                 "Knowledge base articles contributed: 156"]:
        p = body5.add_paragraph()
        p.text = item

    # --- Slide 6: Tools & Resources ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Recommended Tools & Resources"
    body6 = slide6.placeholders[1].text_frame
    body6.text = "Project Management: Jira, Asana, Monday.com"
    for item in ["Communication: Slack, Microsoft Teams",
                 "Documentation: Confluence, Notion, GitBook",
                 "Version Control: Git, GitHub, GitLab",
                 "Design: Figma, Miro for collaborative whiteboarding"]:
        p = body6.add_paragraph()
        p.text = item

    # --- Slide 7: Team Roles ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Defining Team Roles"
    body7 = slide7.placeholders[1].text_frame
    body7.text = "Product Owner: Sarah Chen - Prioritizes backlog"
    for item in ["Scrum Master: David Park - Facilitates ceremonies",
                 "Tech Lead: Priya Sharma - Architecture decisions",
                 "QA Lead: James Wilson - Quality standards",
                 "UX Designer: Maria Garcia - User experience"]:
        p = body7.add_paragraph()
        p.text = item

    # --- Slide 8: Timeline ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Implementation Timeline"
    body8 = slide8.placeholders[1].text_frame
    body8.text = "Week 1-2: Assessment and planning phase"
    for item in ["Week 3-4: Process redesign and tool setup",
                 "Week 5-6: Pilot program with Team Alpha",
                 "Week 7-8: Full rollout and training sessions",
                 "Week 9-10: Monitoring and adjustment period"]:
        p = body8.add_paragraph()
        p.text = item

    # --- Slide 9: Feedback ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "Participant Feedback"
    body9 = slide9.placeholders[1].text_frame
    body9.text = '"The retrospective format changed how our team communicates." - Engineering'
    for item in ['"Finally, a structured approach to documentation." - Product',
                 '"Sprint velocity metrics gave us real visibility." - Management',
                 '"Cross-team collaboration improved dramatically." - Design']:
        p = body9.add_paragraph()
        p.text = item

    # --- Slide 10: Thank You ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[0])
    slide10.shapes.title.text = "Thank You"
    slide10.placeholders[1].text = "Questions? Reach out at workshops@company.com\nNext session: June 15, 2025"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
