"""
Initial Setup: Pie chart on slide 10 shows outdated proportions, not matching the data table.
Task ID: impress_fix_017
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_017'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_slide_with_title(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "FY2026 Budget Allocation Report",
                    "Prepared by the Finance Department | Q1 Review")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "1. Executive Summary",
        "2. Revenue Overview",
        "3. Department Budgets",
        "4. Headcount Planning",
        "5. Capital Expenditure",
        "6. Cost Reduction Initiatives",
        "7. Risk Assessment",
        "8. Q2 Projections",
        "9. Budget Allocation Breakdown",
    ])

    # Slide 3: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Total annual budget: $24.8M (up 12% from FY2025)",
        "Engineering receives largest allocation at 40%",
        "Marketing budget increased to support product launch",
        "Operations streamlined to 10% through automation",
        "Sales team maintained at current capacity",
    ])

    # Slide 4: Revenue Overview
    add_content_slide(prs, "Revenue Overview", [
        "Q1 Revenue: $6.2M (target: $5.9M)",
        "Subscription revenue grew 18% YoY",
        "Enterprise contracts: 34 new signed",
        "Customer retention rate: 94.2%",
        "Average deal size: $182K (up from $156K)",
    ])

    # Slide 5: Department Budgets
    slide5 = add_blank_slide_with_title(prs, "Department Budget Details")
    table_shape = slide5.shapes.add_table(
        6, 4, Inches(1), Inches(1.2), Inches(10), Inches(3.5)
    )
    table = table_shape.table
    headers = ["Department", "FY2025 Budget", "FY2026 Budget", "Change"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)

    dept_data = [
        ["Engineering", "$8,400,000", "$9,920,000", "+18.1%"],
        ["Marketing", "$7,200,000", "$8,680,000", "+20.6%"],
        ["Sales", "$3,100,000", "$3,720,000", "+20.0%"],
        ["Operations", "$2,800,000", "$2,480,000", "-11.4%"],
        ["Total", "$21,500,000", "$24,800,000", "+15.3%"],
    ]
    for r, row_data in enumerate(dept_data, 1):
        for c, val in enumerate(row_data):
            table.cell(r, c).text = val

    # Slide 6: Headcount Planning
    add_content_slide(prs, "Headcount Planning", [
        "Current headcount: 312 employees",
        "Planned new hires: 48 positions",
        "Engineering: 22 new roles (ML, Platform, Security)",
        "Marketing: 14 new roles (Content, Growth, Events)",
        "Sales: 8 new roles (Enterprise, Mid-Market)",
        "Operations: 4 new roles (DevOps, IT Support)",
    ])

    # Slide 7: Capital Expenditure
    add_content_slide(prs, "Capital Expenditure Plan", [
        "Total CapEx: $3.2M",
        "Cloud infrastructure expansion: $1.4M",
        "Office renovation (Building C): $800K",
        "New development hardware: $520K",
        "Security systems upgrade: $480K",
    ])

    # Slide 8: Cost Reduction Initiatives
    add_content_slide(prs, "Cost Reduction Initiatives", [
        "Migrate legacy services to cloud: -$340K/year",
        "Automate QA pipeline: -$180K/year",
        "Renegotiate vendor contracts: -$220K/year",
        "Consolidate office space: -$150K/year",
        "Total projected savings: $890K annually",
    ])

    # Slide 9: Q2 Projections
    add_content_slide(prs, "Q2 Projections", [
        "Revenue target: $7.1M",
        "New product launch expected mid-April",
        "Marketing campaign budget: $1.2M",
        "Expected customer acquisition: 2,800 new accounts",
        "Projected net margin: 22%",
    ])

    # Slide 10: Budget Allocation Breakdown with PIE CHART
    # The chart intentionally shows OUTDATED/WRONG proportions
    # The data table on the slide shows the correct values
    slide10 = add_blank_slide_with_title(prs, "Budget Allocation Breakdown")

    # Add a text box with the "correct" data table (what the chart SHOULD show)
    data_box = slide10.shapes.add_textbox(
        Inches(8.5), Inches(1.5), Inches(4), Inches(3)
    )
    tf = data_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Target Allocation:"
    p.runs[0].font.bold = True
    p.runs[0].font.size = Pt(16)

    allocations = [
        ("Marketing", "35%"),
        ("Engineering", "40%"),
        ("Sales", "15%"),
        ("Operations", "10%"),
    ]
    for dept, pct in allocations:
        p2 = tf.add_paragraph()
        p2.text = f"  {dept}: {pct}"
        p2.runs[0].font.size = Pt(14)

    # Add pie chart with OUTDATED proportions (the "frozen" chart)
    # These are intentionally WRONG - showing old/outdated percentages
    chart_data = CategoryChartData()
    chart_data.categories = ['Marketing', 'Engineering', 'Sales', 'Operations']
    # OLD/WRONG values - the chart appears "frozen" at these outdated numbers
    chart_data.add_series('Budget Share', (25, 30, 25, 20))

    chart_frame = slide10.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.8), Inches(1.5), Inches(7), Inches(5.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Add data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_category_name = True
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.font.size = Pt(11)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
