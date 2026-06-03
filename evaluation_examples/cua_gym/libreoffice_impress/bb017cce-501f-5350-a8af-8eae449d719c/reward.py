"""
Reward Script: Verify pie chart data on slide 10 reflects correct budget allocation
Task ID: impress_fix_017
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Marketing value == 35
  Component 2 (0.25): Engineering value == 40
  Component 3 (0.25): Sales value == 15
  Component 4 (0.25): Operations value == 10
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_fix_017'

# Expected chart data: category -> value
EXPECTED_DATA = {
    'Marketing': 35.0,
    'Engineering': 40.0,
    'Sales': 15.0,
    'Operations': 10.0,
}

def persist_app_state(domain: str):
    """Save any unsaved LibreOffice edits before verification."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify that the pie chart on slide 10 has the correct data values.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 10 slides
    if len(prs.slides) < 10:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 10")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[9]  # slide 10 (0-indexed)

    # Precondition: find a PIE chart on slide 10
    chart = None
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if hasattr(shape, 'chart'):
                    chart = shape.chart
                    break
        except Exception:
            continue

    if chart is None:
        print("FAIL: No chart found on slide 10")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: chart is a PIE-family type (PIE or PIE_EXPLODED)
    try:
        from pptx.enum.chart import XL_CHART_TYPE
        pie_types = {XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED}
        if chart.chart_type not in pie_types:
            print(f"FAIL: Chart is not a PIE type, got {chart.chart_type}")
            print("REWARD: 0.0")
            return 0.0
        print(f"Chart type: {chart.chart_type}")
    except Exception as e:
        print(f"WARN: Could not verify chart type: {e}")

    # Get chart data
    try:
        plot = chart.plots[0]
        categories = list(plot.categories)
        values = list(plot.series[0].values)
        print(f"Chart categories: {categories}")
        print(f"Chart values: {values}")
    except Exception as e:
        print(f"FAIL: Could not read chart data: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Build a mapping of category -> value from the chart
    chart_data = {}
    for cat, val in zip(categories, values):
        chart_data[cat] = val

    # Component 1: Marketing == 35 (0.25 points)
    try:
        actual = chart_data.get('Marketing')
        if actual is not None and abs(actual - EXPECTED_DATA['Marketing']) < 0.5:
            print(f"PASS: Component 1 — Marketing = {actual} (expected 35) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Marketing = {actual}, expected 35")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Engineering == 40 (0.25 points)
    try:
        actual = chart_data.get('Engineering')
        if actual is not None and abs(actual - EXPECTED_DATA['Engineering']) < 0.5:
            print(f"PASS: Component 2 — Engineering = {actual} (expected 40) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 2 — Engineering = {actual}, expected 40")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Sales == 15 (0.25 points)
    try:
        actual = chart_data.get('Sales')
        if actual is not None and abs(actual - EXPECTED_DATA['Sales']) < 0.5:
            print(f"PASS: Component 3 — Sales = {actual} (expected 15) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 3 — Sales = {actual}, expected 15")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Operations == 10 (0.25 points)
    try:
        actual = chart_data.get('Operations')
        if actual is not None and abs(actual - EXPECTED_DATA['Operations']) < 0.5:
            print(f"PASS: Component 4 — Operations = {actual} (expected 10) (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 4 — Operations = {actual}, expected 10")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
