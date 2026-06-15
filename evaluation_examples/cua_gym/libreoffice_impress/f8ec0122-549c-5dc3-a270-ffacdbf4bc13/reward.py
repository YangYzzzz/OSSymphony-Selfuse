"""
FINAL REWARD SCRIPT - SUCCESS
Task: Quick LibreOffice Impress tweak: on slide 72 I need the background switched to the preset “Dark Blue 1” (hex #0B5394) and the title text changed to pure white (#FFFFFF). What’s the fastest way to do that?
Generated: 2025-09-10 17:34:16
Status: success
Model: azure-o3
Total Steps: 3
"""

import os
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import PP_PLACEHOLDER

"""
Reward Script for:
Quick LibreOffice Impress tweak: on slide 72 the background must be preset “Dark Blue 1” (#0B5394) and the title text must be pure white (#FFFFFF).

The script verifies two independent requirements and awards 0.5 points for each:
1. Slide-72 background colour is the exact RGB #0B5394 and uses a solid fill.
2. All explicitly coloured text runs in the slide-72 title placeholder are pure white (#FFFFFF).

The final reward is the sum of both checks (max 1.0).
"""

def verify_background(slide, target_hex: str) -> bool:
    """Verify that slide background is a solid fill with the target_hex colour."""
    try:
        fill = slide.background.fill
        if fill.type != MSO_FILL.SOLID:
            print("✗ Background fill is not SOLID")
            return False
        rgb = fill.fore_color.rgb  # type: ignore[attr-defined]
        if rgb is None:
            print("✗ Background RGB colour not explicitly set")
            return False
        colour_hex = str(rgb).upper()
        if colour_hex == target_hex.upper():
            print(f"✓ Background colour matches target ({colour_hex})")
            return True
        print(f"✗ Background colour {colour_hex} does not match target {target_hex}")
        return False
    except Exception as e:
        print(f"✗ Error verifying background: {e}")
        return False


def verify_title_white(slide) -> bool:
    """Verify that every explicitly coloured run in the title placeholder is white (#FFFFFF)."""
    WHITE = "FFFFFF"
    title_shapes = []

    # Locate title placeholders (TITLE or CENTER_TITLE)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_title = False
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type  # type: ignore[attr-defined]
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                is_title = True
        if shape.name and "Title" in shape.name:
            is_title = True
        if is_title:
            title_shapes.append(shape)

    if not title_shapes:
        print("✗ No title placeholder found on slide 72")
        return False

    any_run_coloured_white = False
    all_coloured_runs_white = True

    for shape in title_shapes:
        for paragraph in shape.text_frame.paragraphs:  # type: ignore[attr-defined]
            for run in paragraph.runs:
                rgb = run.font.color.rgb  # type: ignore[attr-defined]
                if rgb is None:
                    # Run inherits colour – ignore for explicit colour check
                    continue
                colour_hex = str(rgb).upper()
                if colour_hex != WHITE:
                    all_coloured_runs_white = False
                else:
                    any_run_coloured_white = True

    if not any_run_coloured_white:
        print("✗ No text run explicitly set to white in title")
        return False
    if all_coloured_runs_white:
        print("✓ All explicitly coloured runs in title are pure white")
        return True
    print("✗ Some explicitly coloured runs in title are not white")
    return False


def verify_impress_tweak(file_path: str) -> float:
    """Main verification wrapper – returns progressive score between 0.0 and 1.0."""
    print("Starting verification for LibreOffice Impress tweak task…")
    score = 0.0
    TARGET_BG_HEX = "0B5394"

    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Unable to load presentation: {e}")
        return 0.0

    slide_idx = 71  # slide 72 (0-based index)
    if len(prs.slides) <= slide_idx:
        print(f"✗ Presentation only has {len(prs.slides)} slides – slide 72 missing")
        return 0.0

    slide = prs.slides[slide_idx]

    # Check background colour (0.5 pts)
    if verify_background(slide, TARGET_BG_HEX):
        score += 0.5

    # Check title text colour (0.5 pts)
    if verify_title_white(slide):
        score += 0.5

    final_score = min(score, 1.0)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/quick_libreoffice_impress_tweak_on_slide_72_i_need_the_background_switched_to_the_preset_dark_blue_1_golden.pptx"
    reward = verify_impress_tweak(FILE_PATH)
    print(f"REWARD: {reward}")
