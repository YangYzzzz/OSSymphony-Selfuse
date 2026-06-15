"""
Initial Setup: Create a promotional presentation with red background, no Fontwork.
Task ID: impress_ndo_059
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_059'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Promotional hero slide with red background ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Title Only layout
    fill1 = slide1.background.fill
    fill1.solid()
    fill1.fore_color.rgb = RGBColor(0xCC, 0x00, 0x00)  # Red background

    # Add a title textbox at top
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Spring Promotion 2025"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle text
    txBox2 = slide1.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Limited time offers on selected items"
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # --- Slide 2: Product highlights ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    txBox3 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Featured Products"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    products = [
        ("Premium Wireless Headphones", "$149.99", "30% off"),
        ("Smart Fitness Watch", "$89.99", "25% off"),
        ("Portable Bluetooth Speaker", "$59.99", "40% off"),
        ("Noise-Canceling Earbuds", "$119.99", "20% off"),
    ]
    y_pos = Inches(1.5)
    for name, price, discount in products:
        txb = slide2.shapes.add_textbox(Inches(1), y_pos, Inches(7.5), Inches(0.8))
        tf_prod = txb.text_frame
        tf_prod.word_wrap = True
        p_name = tf_prod.paragraphs[0]
        r_name = p_name.add_run()
        r_name.text = f"{name}  —  {price}  "
        r_name.font.name = "Arial"
        r_name.font.size = Pt(18)
        r_name.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        r_disc = p_name.add_run()
        r_disc.text = discount
        r_disc.font.name = "Arial"
        r_disc.font.size = Pt(18)
        r_disc.font.bold = True
        r_disc.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
        y_pos += Inches(1.0)

    # --- Slide 3: Terms and conditions ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    txBox4 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Terms & Conditions"
    p4.alignment = PP_ALIGN.LEFT
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(28)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    terms = [
        "Offer valid from March 15 to April 30, 2025.",
        "Discounts apply to in-store and online purchases.",
        "Cannot be combined with other promotions or coupons.",
        "While supplies last. No rain checks available.",
        "Returns accepted within 30 days with original receipt.",
    ]
    txBox5 = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7.5), Inches(4))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    for i, term in enumerate(terms):
        if i == 0:
            p_t = tf5.paragraphs[0]
        else:
            p_t = tf5.add_paragraph()
        p_t.text = f"• {term}"
        p_t.space_after = Pt(8)
        r_t = p_t.runs[0]
        r_t.font.name = "Arial"
        r_t.font.size = Pt(14)
        r_t.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
