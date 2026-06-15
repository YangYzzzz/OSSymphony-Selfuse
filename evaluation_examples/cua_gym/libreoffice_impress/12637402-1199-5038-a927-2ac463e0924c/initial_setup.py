"""
Initial Setup: Apply animations to slide 9 title and image
Task ID: impress_gf1_047
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

WORKDIR = '/home/user'
TASK_ID = 'impress_gf1_047'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_placeholder_image(path, width=400, height=300, color=(70, 130, 180)):
    """Create a simple product image placeholder."""
    img = Image.new('RGB', (width, height), color)
    # Draw a simple border-like rectangle
    pixels = img.load()
    for x in range(width):
        for y in range(height):
            if x < 5 or x >= width - 5 or y < 5 or y >= height - 5:
                pixels[x, y] = (40, 80, 120)
            elif 100 < x < 300 and 80 < y < 220:
                pixels[x, y] = (220, 220, 240)
    img.save(path)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
    return slide


def add_blank_slide_with_title(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    img_path = f'{WORKDIR}/_product_img.png'
    create_placeholder_image(img_path)

    # --- Slide 1: Title ---
    add_title_slide(prs, "Dynamic Show", "Product Launch Presentation 2025")

    # --- Slide 2: Agenda ---
    add_content_slide(prs, "Agenda", [
        "Company Overview",
        "Market Analysis",
        "Product Features",
        "Technical Architecture",
        "Customer Testimonials",
        "Pricing & Plans",
        "Roadmap",
        "Q&A Session",
    ])

    # --- Slide 3: Company Overview ---
    add_content_slide(prs, "Company Overview", [
        "Founded in 2018 by Sarah Chen and Marcus Johnson",
        "Headquarters in San Francisco, offices in London and Tokyo",
        "Over 500 employees across 12 countries",
        "Revenue growth of 142% year-over-year",
        "Named Top 50 Innovative Companies by FastTrack Magazine",
    ])

    # --- Slide 4: Market Analysis ---
    add_content_slide(prs, "Market Analysis", [
        "Total addressable market: $47.3 billion by 2026",
        "Current market penetration: 8.2%",
        "Primary competitors: AlphaSync, DataBridge, CloudPeak",
        "Key differentiator: Real-time analytics with 99.9% uptime",
        "Growth segment: Mid-market enterprises (200-2000 employees)",
    ])

    # --- Slide 5: Product Features ---
    add_content_slide(prs, "Product Features", [
        "Intelligent dashboard with customizable widgets",
        "AI-powered predictive analytics engine",
        "Seamless integration with 150+ business tools",
        "Role-based access control and audit logging",
        "Real-time collaboration with team annotations",
    ])

    # --- Slide 6: Technical Architecture ---
    add_content_slide(prs, "Technical Architecture", [
        "Microservices architecture on Kubernetes",
        "Event-driven data pipeline with Apache Kafka",
        "PostgreSQL + Redis for persistence and caching",
        "GraphQL API gateway with rate limiting",
        "End-to-end encryption at rest and in transit",
    ])

    # --- Slide 7: Customer Testimonials ---
    add_content_slide(prs, "Customer Testimonials", [
        '"Reduced our reporting time by 73%" — Elena Rodriguez, VP Analytics at Meridian Corp',
        '"The most intuitive platform we\'ve deployed" — James Park, CTO at NovaBridge',
        '"ROI within 3 months of implementation" — Aisha Patel, Director at GlobalReach',
    ])

    # --- Slide 8: Pricing & Plans ---
    add_content_slide(prs, "Pricing & Plans", [
        "Starter: $29/user/month — up to 25 users",
        "Professional: $59/user/month — up to 200 users",
        "Enterprise: Custom pricing — unlimited users",
        "All plans include 24/7 support and onboarding",
        "Annual billing saves 20%",
    ])

    # --- Slide 9: Transition Point (KEY SLIDE) ---
    # This slide has a title textbox and a product image. No animations.
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add title text box
    title_box = slide9.shapes.add_textbox(
        Inches(1.0), Inches(0.5), Inches(8.0), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Transition Point"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2E, 0x3A, 0x87)
    p.alignment = PP_ALIGN.CENTER

    # Add product image
    pic = slide9.shapes.add_picture(
        img_path, Inches(2.5), Inches(2.0), Inches(5.0), Inches(4.0)
    )

    # --- Slide 10: Roadmap ---
    add_content_slide(prs, "2025 Roadmap", [
        "Q1: Launch mobile companion app (iOS & Android)",
        "Q2: AI assistant for natural language queries",
        "Q3: Advanced data governance and compliance tools",
        "Q4: Marketplace for third-party integrations",
    ])

    # --- Slide 11: Team ---
    add_content_slide(prs, "Leadership Team", [
        "Sarah Chen — CEO & Co-Founder",
        "Marcus Johnson — CTO & Co-Founder",
        "Lisa Nakamura — VP of Engineering",
        "David Thompson — VP of Sales",
        "Priya Sharma — Head of Product",
    ])

    # --- Slide 12: Q&A / Closing ---
    slide12 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide12.shapes.add_textbox(
        Inches(2), Inches(2.5), Inches(6), Inches(2)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0x2E, 0x3A, 0x87)
    p2 = tf.add_paragraph()
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Clean up temp image
    if os.path.exists(img_path):
        os.remove(img_path)

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
