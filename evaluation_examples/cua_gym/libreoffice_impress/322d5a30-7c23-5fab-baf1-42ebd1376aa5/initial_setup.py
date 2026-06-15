"""
Initial Setup: Reverse the order of all slides in this presentation.
Task ID: osworld_impress_slide_duplication_reorder_011
Domain: libreoffice_impress

Creates a 5-slide research talk presentation in original order:
  Slide 1: Abstract
  Slide 2: Introduction
  Slide 3: Methods
  Slide 4: Results
  Slide 5: Conclusion
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_011'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content_slide(prs, layout_idx, title_text, content_lines):
    """Add a slide with a title and bullet content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_tf = title_shape.text_frame
    for run in title_tf.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

    # Set content
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)
    return slide


def create_initial():
    prs = Presentation()
    # Use widescreen 16:9 layout
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title (Abstract) ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = title_slide.shapes.title
    title_shape.text = "Abstract"
    for run in title_shape.text_frame.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        "Deep Learning Approaches for Automated Medical Image Analysis:\n"
        "A Comparative Study of CNN and Transformer Architectures\n\n"
        "Dr. Elena Marchetti, Dr. James Okafor, Dr. Priya Nair\n"
        "Department of Computational Medicine, Westfield University\n"
        "Journal of Medical Informatics, Vol. 14, 2025"
    )
    for para in subtitle.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(16)

    # --- Slide 2: Introduction ---
    add_title_content_slide(
        prs,
        layout_idx=1,
        title_text="Introduction",
        content_lines=[
            "Medical imaging generates 3.6 billion studies annually worldwide",
            "Manual analysis is time-consuming and subject to inter-rater variability",
            "Deep learning has achieved radiologist-level accuracy in select tasks",
            "Gap remains: robust generalisation across imaging modalities and institutions",
            "Research objectives:",
            "  - Benchmark CNN vs. Transformer models on multi-modal MRI datasets",
            "  - Investigate transfer learning strategies for low-resource settings",
            "  - Develop ensemble approach combining complementary model strengths",
        ],
    )

    # --- Slide 3: Methods ---
    add_title_content_slide(
        prs,
        layout_idx=1,
        title_text="Methods",
        content_lines=[
            "Dataset: 12,400 anonymised brain MRI scans from 3 tertiary hospitals",
            "  - Split: 70% train / 15% validation / 15% test (patient-level)",
            "  - Modalities: T1w, T2w, FLAIR, DWI",
            "Architectures evaluated:",
            "  - ResNet-50, EfficientNet-B4 (CNN baselines)",
            "  - Swin Transformer-B, ViT-L/16 (attention-based)",
            "  - Hybrid: ConvNeXt-Large (best-of-both)",
            "Training: AdamW optimiser, cosine annealing LR, mixed precision (BF16)",
            "Augmentation: random rotation ±15°, horizontal flip, Gaussian noise σ=0.05",
            "Evaluation metrics: AUC-ROC, F1-macro, Dice coefficient (segmentation)",
        ],
    )

    # --- Slide 4: Results ---
    add_title_content_slide(
        prs,
        layout_idx=1,
        title_text="Results",
        content_lines=[
            "ConvNeXt-Large achieves highest overall AUC-ROC: 0.947 (±0.008)",
            "Swin Transformer-B leads on DWI modality: AUC 0.961",
            "ResNet-50 remains competitive on T1w, AUC 0.928, with 4× lower cost",
            "Ensemble (ConvNeXt + Swin) outperforms all single models: AUC 0.953",
            "Transfer learning from ImageNet reduces required labelled data by 62%",
            "Performance on external site C drops by avg 4.1% — domain shift concern",
            "Inference latency: ResNet-50 38 ms/scan vs. ViT-L/16 212 ms/scan (GPU A100)",
            "Statistical significance confirmed (paired t-test, p < 0.001) for top-3 models",
        ],
    )

    # --- Slide 5: Conclusion ---
    add_title_content_slide(
        prs,
        layout_idx=1,
        title_text="Conclusion",
        content_lines=[
            "Hybrid CNN-Transformer architectures deliver state-of-the-art accuracy",
            "Ensemble approach is recommended when latency constraints permit",
            "Transfer learning substantially reduces labelling burden in clinical settings",
            "Domain shift remains a critical challenge for multi-site deployments",
            "Future work:",
            "  - Federated learning to address data-sharing barriers",
            "  - Self-supervised pre-training on unlabelled radiology archives",
            "  - Prospective clinical trial planned at 5 partner hospitals (2026)",
            "Code and pre-trained weights available at: github.com/westfield-med/dlia",
        ],
    )

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
