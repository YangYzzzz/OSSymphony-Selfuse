"""
Reward Script: Reverse the order of all slides in the presentation
Task ID: osworld_impress_slide_duplication_reorder_011
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 1 is 'Conclusion' (first becomes last, last becomes first) — 0.3 pts
  Component 2: Slide 5 is 'Abstract' (original first slide is now last) — 0.3 pts
  Component 3: Full reversed order is correct — all 5 slides in exact reverse — 0.4 pts
  Total: 1.0
"""

import os

from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_slide_duplication_reorder_011'

# Expected reversed slide order based on task context
EXPECTED_ORDER = ['Conclusion', 'Results', 'Methods', 'Introduction', 'Abstract']


def get_slide_title(slide):
    """Extract the title (first non-empty text paragraph) from a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    return text
    return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: check slide count is correct (5 slides)
    num_slides = len(prs.slides)
    if num_slides != 5:
        print(f"CRITICAL: Expected 5 slides, found {num_slides}. Task cannot be verified.")
        print("REWARD: 0.0")
        return 0.0

    # Collect actual slide titles
    actual_titles = []
    for slide in prs.slides:
        title = get_slide_title(slide)
        actual_titles.append(title)

    print(f"Actual slide order: {actual_titles}")
    print(f"Expected slide order: {EXPECTED_ORDER}")

    # Component 1: Slide 1 is 'Conclusion' (0.3 points)
    # The original last slide should now be first — key indicator of reversal
    try:
        first_title = actual_titles[0]
        if first_title == 'Conclusion':
            print(f"PASS: Component 1 — Slide 1 is 'Conclusion' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Slide 1 expected 'Conclusion', found '{first_title}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 5 is 'Abstract' (0.3 points)
    # The original first slide should now be last — key indicator of reversal
    try:
        last_title = actual_titles[4]
        if last_title == 'Abstract':
            print(f"PASS: Component 2 — Slide 5 is 'Abstract' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 2 — Slide 5 expected 'Abstract', found '{last_title}'")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Full reversed order is correct (0.4 points)
    # Verify all 5 slides are in exact reversed order
    try:
        if actual_titles == EXPECTED_ORDER:
            print(f"PASS: Component 3 — Full slide order is correctly reversed: {actual_titles} (0.4 pts)")
            total_score += 0.4
        else:
            print(f"FAIL: Component 3 — Full order mismatch. Expected {EXPECTED_ORDER}, found {actual_titles}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
