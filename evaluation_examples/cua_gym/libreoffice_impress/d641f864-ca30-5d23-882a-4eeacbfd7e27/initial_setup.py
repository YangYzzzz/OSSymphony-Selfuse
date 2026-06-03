"""
Initial Setup: Build a Chemistry Elements presentation with 6 slides.
Slide 4 ('First Elements') is empty - ready for the agent to add the periodic table grid.
Task ID: impress_teach_059
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_059'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_content(slide, title_text, body_lines):
    """Helper to populate a title+content slide."""
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Chemistry Elements"
    slide1.placeholders[1].text = "An Introduction to the Periodic Table"

    # --- Slide 2: Atomic Structure ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide2, "Atomic Structure", [
        "Atoms consist of protons, neutrons, and electrons",
        "Protons carry positive charge and reside in the nucleus",
        "Neutrons have no charge and stabilize the nucleus",
        "Electrons orbit in quantized energy levels",
        "Atomic number equals the count of protons",
    ])

    # --- Slide 3: Element Classification ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide3, "Element Classification", [
        "Metals: high conductivity, malleable, lustrous",
        "Nonmetals: poor conductors, brittle in solid form",
        "Metalloids: intermediate properties between metals and nonmetals",
        "Noble gases: extremely low reactivity, full valence shells",
        "Transition metals: variable oxidation states, colored compounds",
    ])

    # --- Slide 4: First Elements (EMPTY - only title) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only a title textbox at the top
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "First Elements"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 5: Chemical Bonds ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide5, "Chemical Bonds", [
        "Ionic bonds: transfer of electrons between atoms",
        "Covalent bonds: sharing of electron pairs",
        "Metallic bonds: delocalized electron sea model",
        "Hydrogen bonds: weak intermolecular attraction",
        "Van der Waals forces: temporary dipole interactions",
    ])

    # --- Slide 6: Summary ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_content(slide6, "Summary", [
        "The periodic table organizes elements by atomic number",
        "Element properties follow periodic trends",
        "Understanding atomic structure explains chemical behavior",
        "The first four elements demonstrate fundamental principles",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
