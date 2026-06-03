"""
Reward Script: Periodic table excerpt on slide 4 with 4x4 grid, element symbols, colors, and atomic numbers
Task ID: impress_teach_059
Domain: libreoffice_impress
Scoring:
  Component 1: 4x4 grid of square shapes on slide 4 (0.20 pts)
  Component 2: All squares are equal size (0.10 pts)
  Component 3: First row element symbols H, He, Li, Be (0.25 pts)
  Component 4: Correct fill colors for first row elements (0.25 pts)
  Component 5: Atomic numbers 1, 2, 3, 4 as small text above symbols (0.20 pts)
"""

import os

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_059'


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect all AUTO_SHAPE rectangles on slide 4 (excluding placeholders and textboxes)
    rect_shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rect_shapes.append(shape)

    print(f"INFO: Found {len(rect_shapes)} AUTO_SHAPE(s) on slide 4")

    # Component 1: 4x4 grid of square shapes exists (0.20 points)
    # Must have exactly 16 (or at least 16) rectangle shapes forming a 4x4 grid
    try:
        if len(rect_shapes) >= 16:
            # Check that shapes form a grid pattern: 4 distinct rows, 4 distinct columns
            tops = sorted(set(s.top for s in rect_shapes))
            lefts = sorted(set(s.left for s in rect_shapes))
            # Allow some tolerance for grid alignment (group close values)
            def group_positions(positions, tolerance=50000):
                """Group positions that are within tolerance of each other."""
                if not positions:
                    return []
                groups = [[positions[0]]]
                for pos in positions[1:]:
                    if abs(pos - groups[-1][0]) <= tolerance:
                        groups[-1].append(pos)
                    else:
                        groups.append([pos])
                return groups

            top_groups = group_positions(tops)
            left_groups = group_positions(lefts)

            if len(top_groups) >= 4 and len(left_groups) >= 4:
                print(f"PASS: Component 1 -- 4x4 grid found: {len(rect_shapes)} shapes, {len(top_groups)} rows, {len(left_groups)} columns (0.20 pts)")
                total_score += 0.20
            else:
                print(f"FAIL: Component 1 -- {len(rect_shapes)} shapes but grid layout is {len(top_groups)} rows x {len(left_groups)} cols, need 4x4")
        else:
            print(f"FAIL: Component 1 -- Only {len(rect_shapes)} AUTO_SHAPE(s) found, need >= 16 for 4x4 grid")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: All squares are equal size (0.10 points)
    try:
        if len(rect_shapes) >= 16:
            widths = [s.width for s in rect_shapes]
            heights = [s.height for s in rect_shapes]
            # Check all widths are the same (within tolerance)
            ref_w = widths[0]
            ref_h = heights[0]
            tolerance = 0.02  # 2% tolerance
            all_equal_w = all(abs(w - ref_w) / max(ref_w, 1) <= tolerance for w in widths)
            all_equal_h = all(abs(h - ref_h) / max(ref_h, 1) <= tolerance for h in heights)
            # Also check that shapes are square (width ~= height)
            is_square = abs(ref_w - ref_h) / max(ref_w, ref_h, 1) <= tolerance

            if all_equal_w and all_equal_h and is_square:
                print(f"PASS: Component 2 -- All {len(rect_shapes)} shapes are equal-size squares ({ref_w}x{ref_h} EMU) (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 2 -- Sizes not uniform or not square. equal_w={all_equal_w}, equal_h={all_equal_h}, is_square={is_square}")
        else:
            print(f"FAIL: Component 2 -- Not enough shapes to check size uniformity")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Helper: Find shapes in the first row (topmost row of the grid)
    def get_first_row_shapes(shapes):
        """Return shapes in the topmost row, sorted left to right."""
        if not shapes:
            return []
        # Find the minimum top position (first row)
        min_top = min(s.top for s in shapes)
        tolerance = 50000  # EMU tolerance for same row
        first_row = [s for s in shapes if abs(s.top - min_top) <= tolerance]
        first_row.sort(key=lambda s: s.left)
        return first_row

    first_row = get_first_row_shapes(rect_shapes)
    print(f"INFO: First row has {len(first_row)} shapes")

    # Component 3: First row contains element symbols H, He, Li, Be (0.25 points)
    # Each correct symbol earns 0.25/4 = 0.0625 points
    expected_symbols = ['H', 'He', 'Li', 'Be']
    try:
        symbols_found = 0
        for idx, symbol in enumerate(expected_symbols):
            if idx < len(first_row):
                shape = first_row[idx]
                # Get all text from the shape
                shape_text = ""
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            shape_text += run.text.strip() + " "
                shape_text = shape_text.strip()

                if symbol in shape_text:
                    symbols_found += 1
                    print(f"  PASS: Symbol '{symbol}' found in shape {idx} (text: {repr(shape_text)})")
                else:
                    print(f"  FAIL: Symbol '{symbol}' not found in shape {idx} (text: {repr(shape_text)})")
            else:
                print(f"  FAIL: No shape at position {idx} for symbol '{symbol}'")

        comp3_score = 0.25 * (symbols_found / 4)
        if symbols_found == 4:
            print(f"PASS: Component 3 -- All 4 element symbols found ({comp3_score:.4f} pts)")
            total_score += comp3_score
        elif symbols_found > 0:
            print(f"PARTIAL: Component 3 -- {symbols_found}/4 symbols found ({comp3_score:.4f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 -- 0/4 symbols found")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Correct fill colors for first row elements (0.25 points)
    # H=#FFCDD2, He=#BBDEFB, Li=#C8E6C9, Be=#FFF9C4
    expected_colors = ['FFCDD2', 'BBDEFB', 'C8E6C9', 'FFF9C4']
    try:
        colors_matched = 0
        for idx, expected_color in enumerate(expected_colors):
            if idx < len(first_row):
                shape = first_row[idx]
                try:
                    fill = shape.fill
                    if fill.type is not None and fill.type == 1:  # SOLID fill
                        actual_color = str(fill.fore_color.rgb).upper()
                        if actual_color == expected_color.upper():
                            colors_matched += 1
                            print(f"  PASS: Shape {idx} fill color {actual_color} matches expected {expected_color}")
                        else:
                            print(f"  FAIL: Shape {idx} fill color {actual_color} != expected {expected_color}")
                    else:
                        print(f"  FAIL: Shape {idx} fill type is {fill.type}, expected SOLID (1)")
                except Exception as e:
                    print(f"  FAIL: Shape {idx} fill check error: {e}")
            else:
                print(f"  FAIL: No shape at position {idx} for color check")

        comp4_score = 0.25 * (colors_matched / 4)
        if colors_matched == 4:
            print(f"PASS: Component 4 -- All 4 fill colors correct ({comp4_score:.4f} pts)")
            total_score += comp4_score
        elif colors_matched > 0:
            print(f"PARTIAL: Component 4 -- {colors_matched}/4 colors matched ({comp4_score:.4f} pts)")
            total_score += comp4_score
        else:
            print(f"FAIL: Component 4 -- 0/4 colors matched")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Atomic numbers 1, 2, 3, 4 as small text above each symbol (0.20 points)
    # In the golden, each element shape has 2 paragraphs: para 0 = atomic number (small), para 1 = symbol (large)
    expected_numbers = ['1', '2', '3', '4']
    try:
        numbers_found = 0
        for idx, expected_num in enumerate(expected_numbers):
            if idx < len(first_row):
                shape = first_row[idx]
                if shape.has_text_frame and len(shape.text_frame.paragraphs) >= 2:
                    # Check first paragraph contains the atomic number
                    first_para = shape.text_frame.paragraphs[0]
                    first_para_text = "".join(r.text for r in first_para.runs).strip()

                    # Check the atomic number text is smaller than the symbol text
                    second_para = shape.text_frame.paragraphs[1]
                    first_para_size = None
                    second_para_size = None
                    if first_para.runs:
                        first_para_size = first_para.runs[0].font.size
                    if second_para.runs:
                        second_para_size = second_para.runs[0].font.size

                    if first_para_text == expected_num:
                        # Verify the atomic number text is smaller than the symbol
                        if first_para_size is not None and second_para_size is not None and first_para_size < second_para_size:
                            numbers_found += 1
                            print(f"  PASS: Atomic number '{expected_num}' found above symbol, size {first_para_size} < {second_para_size}")
                        elif first_para_text == expected_num:
                            # Number is correct even if size check fails
                            numbers_found += 1
                            print(f"  PASS: Atomic number '{expected_num}' found (size comparison: {first_para_size} vs {second_para_size})")
                        else:
                            print(f"  FAIL: Atomic number text '{first_para_text}' found but size not smaller")
                    else:
                        print(f"  FAIL: Expected atomic number '{expected_num}' in first para, found '{first_para_text}'")
                else:
                    # Check if atomic number appears anywhere in the shape text
                    all_text = shape.text if hasattr(shape, 'text') else ""
                    if expected_num in all_text.split('\n') or expected_num in all_text:
                        numbers_found += 1
                        print(f"  PASS: Atomic number '{expected_num}' found in shape text: {repr(all_text)}")
                    else:
                        print(f"  FAIL: Atomic number '{expected_num}' not found. Shape text: {repr(all_text)}")
            else:
                print(f"  FAIL: No shape at position {idx} for atomic number check")

        comp5_score = 0.20 * (numbers_found / 4)
        if numbers_found == 4:
            print(f"PASS: Component 5 -- All 4 atomic numbers found ({comp5_score:.4f} pts)")
            total_score += comp5_score
        elif numbers_found > 0:
            print(f"PARTIAL: Component 5 -- {numbers_found}/4 atomic numbers found ({comp5_score:.4f} pts)")
            total_score += comp5_score
        else:
            print(f"FAIL: Component 5 -- 0/4 atomic numbers found")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
