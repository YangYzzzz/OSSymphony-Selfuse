"""
FINAL REWARD SCRIPT - SUCCESS
Task: Slide 43 feels cluttered, so I’d like to strip it back to just the heading. How do I change that specific slide’s layout to “Title Only” in LibreOffice Impress?
Generated: 2025-09-10 15:16:52
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def verify_slide_title_only(file_path: str, slide_number: int = 43) -> float:
    """Verify that a specific slide in a PPTX has been changed to the
    'Title Only' layout (i.e., it contains only a title placeholder and no
    other body/content placeholders or extra text shapes).

    Parameters
    ----------
    file_path : str
        Path to the PowerPoint (.pptx) file to verify.
    slide_number : int, optional
        Human-friendly slide number to check (defaults to 43).

    Returns
    -------
    float
        A progressive score between 0.0 and 1.0 reflecting the level of
        task completion. Exactly 1.0 means perfect completion.
    """

    print(f"Verifying slide {slide_number} layout ===> 'Title Only'\nFile: {file_path}\n")

    # Initialise scoring
    total_score = 0.0
    max_score = 1.0

    # ------------------------------------------------------------------
    # 1. Basic file existence & loading (prerequisite – NO POINTS)
    # ------------------------------------------------------------------
    if not os.path.exists(file_path):
        print("✗ File not found. Verification aborted.")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        return 0.0

    # ------------------------------------------------------------------
    # 2. Validate slide index exists (prerequisite – NO POINTS)
    # ------------------------------------------------------------------
    index = slide_number - 1  # convert to 0-based index
    if index < 0 or index >= len(prs.slides):
        print(f"✗ Slide {slide_number} does not exist (total slides: {len(prs.slides)})")
        return 0.0

    slide = prs.slides[index]

    # ------------------------------------------------------------------
    # 3. Requirement A – Layout must be exactly 'Title Only'  (0.5 pts)
    # ------------------------------------------------------------------
    layout_name = slide.slide_layout.name if slide.slide_layout else "<None>"
    print(f"Found layout name: '{layout_name}'")

    if layout_name.strip().lower() == "title only":
        total_score += 0.5
        print("✓ Layout correctly set to 'Title Only' (+0.5)")
    else:
        print("✗ Layout is NOT 'Title Only' (0 pts)")

    # ------------------------------------------------------------------
    # 4. Requirement B – Only one TITLE placeholder, no BODY placeholders
    #                      or other unexpected placeholders          (0.3)
    # ------------------------------------------------------------------
    placeholders = [sh for sh in slide.shapes if getattr(sh, "is_placeholder", False)]
    title_ph   = [ph for ph in placeholders if ph.placeholder_format.type == 1]  # TITLE
    body_ph    = [ph for ph in placeholders if ph.placeholder_format.type == 2]  # BODY
    other_ph   = [ph for ph in placeholders if ph.placeholder_format.type not in (1, 2)]

    print(
        f"Placeholder summary  – Total: {len(placeholders)}, "
        f"Title: {len(title_ph)}, Body: {len(body_ph)}, Other: {len(other_ph)}"
    )

    if len(title_ph) == 1 and len(body_ph) == 0 and len(other_ph) == 0:
        total_score += 0.3
        print("✓ Placeholder configuration correct (+0.3)")
    else:
        print("✗ Unexpected placeholders detected (0 pts)")

    # ------------------------------------------------------------------
    # 5. Requirement C – No additional text-bearing shapes besides title
    #                      (guards against rogue text boxes)         (0.2)
    # ------------------------------------------------------------------
    extraneous_texts = []
    for sh in slide.shapes:
        if hasattr(sh, "text") and sh.text and sh.text.strip():
            # Ignore the title placeholder itself
            if not (getattr(sh, "is_placeholder", False) and sh.placeholder_format.type == 1):
                extraneous_texts.append(sh.text.strip())

    if not extraneous_texts:
        total_score += 0.2
        print("✓ No extraneous text shapes found (+0.2)")
    else:
        print("✗ Extraneous text detected (0 pts):")
        for txt in extraneous_texts:
            print(f"   - '{txt[:40]}'")

    # ------------------------------------------------------------------
    # 6. Finalise scoring
    # ------------------------------------------------------------------
    final_score = min(total_score, max_score)
    print(f"\nTotal Score: {final_score}/{max_score}")
    print(f"REWARD: {final_score}")
    return final_score


# -------------------
# Execute verification
# -------------------
if __name__ == "__main__":
    FILE_PATH = "/home/user/slide_43_feels_cluttered_so_id_like_to_strip_it_back_to_just_the_heading_how_do_i_change_that_specif_golden.pptx"
    verify_slide_title_only(FILE_PATH)

