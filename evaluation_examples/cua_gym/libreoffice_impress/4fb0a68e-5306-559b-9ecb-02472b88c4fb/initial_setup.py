"""
Initial Setup: Create a presentation with a multi-series bar chart on slide 2
with legend on the right side.
Task ID: impress_tct_043
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_043'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Quarterly Revenue Comparison"
    slide1.placeholders[1].text = "FY2025 Performance Review — Regional Breakdown"

    # --- Slide 2: Chart Slide (multi-series bar chart, legend on RIGHT) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Revenue by Region (Q1–Q4)"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    # Chart data
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('North America', (42500, 48200, 51800, 55300))
    chart_data.add_series('Europe', (31200, 34800, 37500, 39100))
    chart_data.add_series('Asia Pacific', (28700, 33400, 38900, 44200))
    chart_data.add_series('Latin America', (12800, 14500, 16200, 18900))

    # Add chart — position and size
    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.2),
        Inches(11.5), Inches(5.8),
        chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    # Style the chart
    chart.has_title = False

    # --- Slide 3: Analysis slide ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])

    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf3t = txBox3_title.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Key Insights"
    p3t.alignment = PP_ALIGN.LEFT
    r3t = p3t.runs[0]
    r3t.font.size = Pt(28)
    r3t.font.bold = True
    r3t.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11), Inches(5.5))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    insights = [
        "Asia Pacific demonstrated the strongest growth trajectory at 54% year-over-year, "
        "surpassing Europe in Q3 and narrowing the gap with North America.",
        "North America maintained steady growth of 30% across all quarters, driven primarily "
        "by enterprise SaaS subscriptions and cloud infrastructure services.",
        "Latin America, while the smallest region, showed consistent 48% growth and is "
        "projected to exceed $20K monthly revenue by Q1 2026.",
        "Europe's growth moderated to 25% in the second half due to regulatory changes "
        "affecting data processing services in key markets.",
    ]
    for i, text in enumerate(insights):
        if i == 0:
            p = tf3.paragraphs[0]
        else:
            p = tf3.add_paragraph()
        p.text = text
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # --- Slide 4: Summary slide ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])

    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.7))
    tf4t = txBox4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Summary & Next Steps"
    p4t.alignment = PP_ALIGN.LEFT
    r4t = p4t.runs[0]
    r4t.font.size = Pt(28)
    r4t.font.bold = True
    r4t.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11), Inches(5.5))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True

    summaries = [
        "Total global revenue reached $181.7K in FY2025 Q4, a 37% increase over Q1 baseline.",
        "Recommendation: Increase investment in Asia Pacific sales team and marketing presence.",
        "Action item: Schedule regional strategy sessions with country managers by end of January.",
        "Next review: FY2026 Q1 results presentation scheduled for April 15, 2026.",
    ]
    for i, text in enumerate(summaries):
        if i == 0:
            p = tf4.paragraphs[0]
        else:
            p = tf4.add_paragraph()
        p.text = text
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
