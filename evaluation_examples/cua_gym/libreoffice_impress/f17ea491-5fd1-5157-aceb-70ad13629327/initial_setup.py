"""
Initial Setup: UX Research Presentation (6 slides, no notes, white backgrounds)
Task ID: osworld_impress_note_bg_combined_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_white_background(slide):
    """Set explicit white solid background on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_content_slide(prs, layout_idx, title_text, content_lines):
    """Add a slide with a title and bullet content."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Set content placeholder if available
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(content_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
            break
    return slide


def create_initial():
    prs = Presentation()
    # Standard widescreen: 10 x 7.5 inches
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = add_title_content_slide(prs, 0,
        "UX Research Findings 2025",
        ["A Comprehensive Analysis of User Behavior and Pain Points"]
    )
    set_white_background(slide1)
    # NO notes on slide 1 (task requires adding notes)

    # --- Slide 2: Agenda ---
    slide2 = add_title_content_slide(prs, 1,
        "Today's Agenda",
        [
            "1. Research Overview & Methodology",
            "2. Key Findings from User Interviews",
            "3. Deep Dive: User Research Data",
            "4. Personas & Journey Maps",
            "5. Recommendations & Next Steps",
            "6. Q&A Session",
        ]
    )
    # Slide 2 gets default white background (unchanged by task)

    # --- Slide 3: Research Methodology ---
    slide3 = add_title_content_slide(prs, 1,
        "Research Methodology",
        [
            "48 in-depth user interviews conducted (Feb–Mar 2025)",
            "Participants: 18–55 years, diverse backgrounds",
            "Qualitative coding: 12 core theme categories",
            "Quantitative survey: 1,240 respondents",
            "Usability testing: 5-second tests + think-aloud sessions",
        ]
    )
    set_white_background(slide3)
    # NO notes on slide 3 (task requires adding notes)

    # --- Slide 4: Key Findings ---
    slide4 = add_title_content_slide(prs, 1,
        "Key Findings",
        [
            "76% of users struggled with navigation on first visit",
            "Top pain point: information overload on landing page",
            "Mobile experience rated 3.2/5 vs desktop 4.1/5",
            "42% abandoned checkout due to unclear pricing",
            "Trust indicators significantly increased conversion rate (+28%)",
        ]
    )
    # Slide 4: default background (unchanged)

    # --- Slide 5: Persona Highlights ---
    slide5 = add_title_content_slide(prs, 1,
        "Persona Highlights",
        [
            "Primary persona: 'Tech-Savvy Millennial' (age 28–38)",
            "Secondary persona: 'Busy Professional' (age 39–52)",
            "Goal: quick task completion with minimal friction",
            "Frustration driver: unexpected page redirects",
            "Motivation: transparent pricing and clear value proposition",
        ]
    )
    set_white_background(slide5)
    # NO notes on slide 5 (task requires adding notes)

    # --- Slide 6: Q&A / Closing ---
    slide6 = add_title_content_slide(prs, 1,
        "Q&A and Discussion",
        [
            "Thank you for your time and attention",
            "Full report available on the shared drive",
            "Contact: ux-research@company.com",
            "Follow-up workshop scheduled for April 10, 2025",
        ]
    )
    # Slide 6: default background (unchanged)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
