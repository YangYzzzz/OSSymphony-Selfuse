"""
Reward Script: Apply background and notes to three slides in a presentation
Task ID: osworld_impress_note_bg_combined_008
Domain: libreoffice_impress
Scoring:
  - Component 1: Slide 1 has note "Welcome and housekeeping" (0.20 pts)
  - Component 2: Slide 3 has light yellow background (FFFFE0) (0.20 pts)
  - Component 3: Slide 3 has note "Deep dive into user research data" (0.20 pts)
  - Component 4: Slide 5 has pale green background (E0FFE0) (0.20 pts)
  - Component 5: Slide 5 has note "Call to action and next steps" (0.20 pts)
  Total: 1.0

Notes:
  - Slide 1 background (white/FFFFFF) is already white in initial state and is
    NOT scored (precondition, not a task-introduced change).
  - All scored components FAIL on initial_env and PASS on golden_env.
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_note_bg_combined_008'


def get_slide_notes(slide):
    """Return stripped notes text for a slide, or empty string if none."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def get_slide_bg_rgb(slide):
    """Return background RGB string (e.g., 'FFFFE0') if solid fill, else None."""
    try:
        fill = slide.background.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb)
        return None
    except Exception:
        return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Load presentation — if this fails, nothing can be verified
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: confirm 6 slides exist
    if len(prs.slides) < 6:
        print(f"CRITICAL: Expected at least 6 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide1 = prs.slides[0]
    slide3 = prs.slides[2]
    slide5 = prs.slides[4]

    # Component 1: Slide 1 note = "Welcome and housekeeping" (0.20 pts)
    # Slide 1 background is already white (precondition), so only the note is scored.
    try:
        notes1 = get_slide_notes(slide1)
        expected_notes1 = "Welcome and housekeeping"
        if notes1 == expected_notes1:
            print(f"PASS: Component 1 — Slide 1 note = {repr(notes1)} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Slide 1 note: expected {repr(expected_notes1)}, found {repr(notes1)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 3 background = light yellow (FFFFE0) (0.20 pts)
    try:
        bg3 = get_slide_bg_rgb(slide3)
        expected_bg3 = "FFFFE0"
        if bg3 == expected_bg3:
            print(f"PASS: Component 2 — Slide 3 background = #{bg3} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 2 — Slide 3 background: expected #{expected_bg3}, found {repr(bg3)}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 note = "Deep dive into user research data" (0.20 pts)
    try:
        notes3 = get_slide_notes(slide3)
        expected_notes3 = "Deep dive into user research data"
        if notes3 == expected_notes3:
            print(f"PASS: Component 3 — Slide 3 note = {repr(notes3)} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 3 — Slide 3 note: expected {repr(expected_notes3)}, found {repr(notes3)}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 5 background = pale green (E0FFE0) (0.20 pts)
    try:
        bg5 = get_slide_bg_rgb(slide5)
        expected_bg5 = "E0FFE0"
        if bg5 == expected_bg5:
            print(f"PASS: Component 4 — Slide 5 background = #{bg5} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — Slide 5 background: expected #{expected_bg5}, found {repr(bg5)}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Slide 5 note = "Call to action and next steps" (0.20 pts)
    try:
        notes5 = get_slide_notes(slide5)
        expected_notes5 = "Call to action and next steps"
        if notes5 == expected_notes5:
            print(f"PASS: Component 5 — Slide 5 note = {repr(notes5)} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 5 — Slide 5 note: expected {repr(expected_notes5)}, found {repr(notes5)}")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score:.2f}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point: test against canonical artifact path on VM
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
