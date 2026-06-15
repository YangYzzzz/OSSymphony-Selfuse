"""
FINAL REWARD SCRIPT - SUCCESS
Task: While doing a final proofread of my 120-slide deck, I noticed slide 96 still says “market outlook q3” in sentence case. All the other headings follow Title Case. In LibreOffice Impress, what’s the quickest way to fix that one slide so the title reads “Market Outlook Q3” (capitalize each word)?
Generated: 2025-09-10 23:20:08
Status: success
Model: azure-o3
Total Steps: 1
"""

import os
from pptx import Presentation


def verify_slide_96_title(file_path: str) -> float:
    """Verify that slide 96’s title reads exactly "Market Outlook Q3" in Title Case.

    Scoring (progressive):
    • 0.4 pts – Title text differs from the original all-lowercase version
    • 0.6 pts – Title text matches *exactly* "Market Outlook Q3"

    Returns a float between 0.0 and 1.0 (1.0 = full credit).
    """

    expected_exact = "Market Outlook Q3"
    original_lower = "market outlook q3"  # what needed fixing

    max_score = 1.0
    score = 0.0

    print(f"Loading presentation: {file_path}")

    # 1) File must exist
    if not os.path.exists(file_path):
        print("✗ File does not exist")
        return 0.0

    # 2) Load the presentation
    try:
        prs = Presentation(file_path)
        print(f"✓ Presentation loaded with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Error loading presentation: {e}")
        return 0.0

    # 3) Ensure slide 96 exists (index 95)
    if len(prs.slides) < 96:
        print("✗ Presentation has fewer than 96 slides – task incomplete")
        return 0.0

    slide = prs.slides[95]  # zero-based index

    # 4) Extract title text (prefer the designated title shape, else first text shape)
    title_text = None
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None and title_shape.has_text_frame:
        title_text = title_shape.text.strip()
    else:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    title_text = text
                    break

    if title_text is None:
        print("✗ No text found on slide 96 to verify title")
        return 0.0

    # Normalise whitespace for comparison
    title_text_collapsed = " ".join(title_text.split())
    print(f"Found title text on slide 96: '{title_text_collapsed}'")

    # 5) Scoring – has text changed from original lower-case? (0.4 pts)
    if title_text_collapsed != original_lower:
        score += 0.4
        print("✓ Title text has changed from original lowercase (0.4 pts)")
    else:
        print("✗ Title text is still lowercase (0 pts)")

    # 6) Exact Title Case match? (+0.6 pts)
    if title_text_collapsed == expected_exact:
        score += 0.6
        print("✓ Title text matches exact expected Title Case (0.6 pts)")
    else:
        print("✗ Title text does not exactly match expected Title Case (0 pts)")

    final_score = min(score, max_score)
    print(f"Total score: {final_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/while_doing_a_final_proofread_of_my_120_slide_deck_i_noticed_slide_96_still_says_market_outlook_q3_i_golden.pptx"
    reward = verify_slide_96_title(FILE_PATH)
    print(f"REWARD: {reward}")

