"""
FINAL REWARD SCRIPT - SUCCESS
Task: Insert the author field in the footer, left-aligned.
Generated: 2025-10-17 06:55:33
Status: success
Model: azure-o3
Total Steps: 7
"""

"""
Reward Verification Script for:
Task: Insert the author field in the footer, left-aligned.

This script awards a progressive score based on REAL checks only
(no bias for natural conditions).  A perfect implementation earns 1.0.

Scoring rubric (adds up to 1.0):
1. Author text appears on ≥1 slide ..................................... 0.3
2. Author text appears on *every* slide ................................. 0.2
3. Author text is positioned in footer region (bottom 30 %) ............. 0.3
4. Author text is left-aligned .......................................... 0.2

If any prerequisite fails (file missing, cannot load, no author metadata)
verification stops and returns current score (possibly 0.0).
"""
import os
import zipfile
import lxml.etree as ET
from pptx import Presentation
from pptx.enum.text import PP_ALIGN


def verify_author_footer_left(file_path: str) -> float:
    """Return a progressive score between 0.0 and 1.0 for task completion."""
    max_score = 1.0
    score = 0.0

    # ---------- Prerequisite checks (no points awarded) ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded presentation with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Could not open PPTX: {e}")
        return 0.0

    # ---------- Retrieve author from core properties ----------
    author = None
    try:
        with zipfile.ZipFile(file_path) as pptx_zip:
            core_xml = pptx_zip.read('docProps/core.xml')
            core_root = ET.fromstring(core_xml)
            ns = {
                'cp': 'http://schemas.openxmlformats.org/package/2006/metadata/core-properties',
                'dc': 'http://purl.org/dc/elements/1.1/'
            }
            author_elem = core_root.find('dc:creator', ns)
            if author_elem is not None and author_elem.text:
                author = author_elem.text.strip()
    except Exception as e:
        print(f"✗ Error reading core properties: {e}")

    if not author:
        print("✗ No author metadata found; cannot verify dynamic author field")
        return 0.0

    print(f"Author from metadata: '{author}'")

    # ---------- Slide-by-slide verification ----------
    found_author = []          # per slide: True if author text detected
    left_aligned_ok = []       # per slide: alignment correct when detected
    footer_region_ok = []      # per slide: vertical position correct when detected

    footer_threshold_ratio = 0.7  # >70 % of slide height considered footer

    for i, slide in enumerate(prs.slides):
        slide_found = False
        slide_left_ok = False
        slide_footer_ok = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if author.lower() in shape.text.lower():
                slide_found = True

                # Alignment check (None == default left for placeholders)
                para = shape.text_frame.paragraphs[0]
                slide_left_ok = (para.alignment in (None, PP_ALIGN.LEFT))

                # Footer region check (bottom 30 % of slide)
                slide_height = prs.slide_height
                slide_footer_ok = shape.top > slide_height * footer_threshold_ratio
                break  # assume one author field per slide is sufficient

        found_author.append(slide_found)
        left_aligned_ok.append(slide_left_ok if slide_found else False)
        footer_region_ok.append(slide_footer_ok if slide_found else False)

        print(
            f"Slide {i+1}: found={slide_found}, "
            f"leftAligned={slide_left_ok}, footerRegion={slide_footer_ok}"
        )

    total_slides = len(prs.slides)
    slides_with_author = sum(found_author)

    # ---------- Scoring ----------
    if slides_with_author > 0:
        score += 0.3
        print("✓ Author text found on at least one slide (+0.3)")
    else:
        print("✗ Author text not found on any slide (0 pts)")
        return score  # cannot earn further points without this

    if slides_with_author == total_slides and total_slides > 0:
        score += 0.2
        print("✓ Author text present on every slide (+0.2)")
    else:
        print("✗ Author text missing from some slides (+0 pts)")

    if all((not f) or ok for f, ok in zip(found_author, footer_region_ok)):
        score += 0.3
        print("✓ Author text positioned in footer region (+0.3)")
    else:
        print("✗ Author text not consistently in footer region (+0 pts)")

    if all((not f) or ok for f, ok in zip(found_author, left_aligned_ok)):
        score += 0.2
        print("✓ Author text left-aligned (+0.2)")
    else:
        print("✗ Author text not consistently left-aligned (+0 pts)")

    final_score = min(score, max_score)
    print(f"Total score: {final_score:.1f} / {max_score}")
    return final_score


if __name__ == "__main__":
    FILE_PATH = "/home/user/insert_the_author_field_in_the_footer_left_aligned.pptx"
    reward = verify_author_footer_left(FILE_PATH)
    print(f"REWARD: {reward}")

