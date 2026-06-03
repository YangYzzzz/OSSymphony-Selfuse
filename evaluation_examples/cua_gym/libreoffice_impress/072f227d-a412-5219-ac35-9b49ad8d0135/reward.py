"""
Reward Script: Create a stacked bar chart on slide 3 with two data series
Task ID: impress_tct_045
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) — Chart exists on slide 3 and is stacked bar type
  Component 2 (0.30) — Two series named 'Online Sales' and 'In-Store Sales'
  Component 3 (0.20) — Four categories (North, South, East, West)
  Component 4 (0.20) — Legend is present
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_045'

# XML namespaces for chart parsing
NS_C = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def find_chart_on_slide3(pptx_path):
    """
    Find chart XML for slide 3 by following relationships.
    Returns parsed chart XML root or None.
    """
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        # Read slide3.xml relationships to find chart reference
        slide3_rels_path = 'ppt/slides/_rels/slide3.xml.rels'
        try:
            rels_xml = zf.read(slide3_rels_path).decode('utf-8')
        except KeyError:
            print("FAIL: slide3 relationships file not found")
            return None

        rels_root = ET.fromstring(rels_xml)
        chart_target = None
        for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
            rel_type = rel.get('Type', '')
            if 'chart' in rel_type.lower():
                chart_target = rel.get('Target', '')
                break

        if not chart_target:
            print("FAIL: No chart relationship found in slide 3")
            return None

        # Resolve relative path (e.g., ../charts/chart1.xml)
        chart_path = 'ppt/charts/' + chart_target.split('/')[-1]
        try:
            chart_xml = zf.read(chart_path).decode('utf-8')
            return ET.fromstring(chart_xml)
        except KeyError:
            print(f"FAIL: Chart file not found at {chart_path}")
            return None


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: file must exist and have at least 3 slides
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        if len(prs.slides) < 3:
            print(f"CRITICAL: Presentation has only {len(prs.slides)} slides, need at least 3")
            print("REWARD: 0.0")
            return 0.0
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Chart exists on slide 3 and is stacked bar type (0.30 points)
    try:
        chart_root = find_chart_on_slide3(file_path)
        if chart_root is None:
            print("FAIL: Component 1 — No chart found on slide 3 (0.0 pts)")
        else:
            # Check for barChart element with stacked grouping
            bar_chart = chart_root.find(f'.//{{{NS_C}}}barChart')
            if bar_chart is None:
                print("FAIL: Component 1 — Chart found but not a bar chart type")
            else:
                grouping_el = bar_chart.find(f'{{{NS_C}}}grouping')
                grouping_val = grouping_el.get('val', '') if grouping_el is not None else ''
                if grouping_val == 'stacked':
                    print(f"PASS: Component 1 — Stacked bar chart found on slide 3 (0.30 pts)")
                    total_score += 0.30
                else:
                    print(f"FAIL: Component 1 — Bar chart grouping is '{grouping_val}', expected 'stacked'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Two series named 'Online Sales' and 'In-Store Sales' (0.30 points)
    try:
        if chart_root is not None:
            bar_chart = chart_root.find(f'.//{{{NS_C}}}barChart')
            series_elements = bar_chart.findall(f'{{{NS_C}}}ser') if bar_chart is not None else []
            series_names = []
            for ser in series_elements:
                tx = ser.find(f'{{{NS_C}}}tx')
                if tx is not None:
                    # Look for cached string value
                    v_el = tx.find(f'.//{{{NS_C}}}v')
                    if v_el is not None and v_el.text:
                        series_names.append(v_el.text.strip())

            print(f"  Found series names: {series_names}")
            expected_names = {'Online Sales', 'In-Store Sales'}
            found_names = set(series_names)

            if len(series_elements) == 2 and expected_names == found_names:
                print(f"PASS: Component 2 — Two series 'Online Sales' and 'In-Store Sales' found (0.30 pts)")
                total_score += 0.30
            elif len(series_elements) == 2:
                # Partial: right count but wrong names
                matching = expected_names.intersection(found_names)
                if len(matching) >= 1:
                    print(f"PARTIAL: Component 2 — 2 series but only matched names: {matching} (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component 2 — 2 series but names don't match: {found_names}")
            elif len(series_elements) >= 1:
                print(f"FAIL: Component 2 — Expected 2 series, found {len(series_elements)}")
            else:
                print("FAIL: Component 2 — No series found in chart")
        else:
            print("FAIL: Component 2 — No chart to check series")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Four categories (North, South, East, West) (0.20 points)
    try:
        if chart_root is not None:
            bar_chart = chart_root.find(f'.//{{{NS_C}}}barChart')
            categories = []
            if bar_chart is not None:
                # Categories are in the first series cat element
                first_ser = bar_chart.find(f'{{{NS_C}}}ser')
                if first_ser is not None:
                    cat_el = first_ser.find(f'{{{NS_C}}}cat')
                    if cat_el is not None:
                        for pt in cat_el.findall(f'.//{{{NS_C}}}pt'):
                            v_el = pt.find(f'{{{NS_C}}}v')
                            if v_el is not None and v_el.text:
                                categories.append(v_el.text.strip())

            print(f"  Found categories: {categories}")
            expected_cats = {'North', 'South', 'East', 'West'}
            found_cats = set(categories)

            if len(categories) == 4 and expected_cats == found_cats:
                print(f"PASS: Component 3 — Four region categories found (0.20 pts)")
                total_score += 0.20
            elif len(found_cats.intersection(expected_cats)) >= 2:
                print(f"PARTIAL: Component 3 — Some categories match: {found_cats} (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 3 — Categories don't match expected regions: {categories}")
        else:
            print("FAIL: Component 3 — No chart to check categories")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Legend is present (0.20 points)
    try:
        if chart_root is not None:
            chart_el = chart_root.find(f'{{{NS_C}}}chart')
            legend_el = chart_el.find(f'{{{NS_C}}}legend') if chart_el is not None else None
            if legend_el is not None:
                print(f"PASS: Component 4 — Legend present in chart (0.20 pts)")
                total_score += 0.20
            else:
                print("FAIL: Component 4 — No legend found in chart")
        else:
            print("FAIL: Component 4 — No chart to check legend")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
