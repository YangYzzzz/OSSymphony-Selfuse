"""
Initial Setup: Create a 5-slide sales presentation with no chart on slide 3.
Task ID: impress_tct_045
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tct_045'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Sales Channels Report Q1 2025"
    slide1.placeholders[1].text = "Prepared by Analytics Division\nApril 2025"

    # --- Slide 2: Executive Summary ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Executive Summary"
    body2 = slide2.placeholders[1].text_frame
    body2.text = "Total revenue across all channels reached $4.8M in Q1 2025"
    p2a = body2.add_paragraph()
    p2a.text = "Online sales grew 18% year-over-year, driven by mobile commerce"
    p2a.level = 0
    p2b = body2.add_paragraph()
    p2b.text = "In-store sales remained steady with a 3% increase in foot traffic"
    p2b.level = 0
    p2c = body2.add_paragraph()
    p2c.text = "The North region continues to lead in both channels combined"
    p2c.level = 0

    # --- Slide 3: Sales by Channel and Region (NO chart - task asks agent to create it) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Sales by Channel and Region"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x2E, 0x4A, 0x7A)

    # Add a note text box indicating data is available
    note_box = slide3.shapes.add_textbox(Inches(1.5), Inches(3), Inches(6), Inches(1))
    ntf = note_box.text_frame
    np = ntf.paragraphs[0]
    np.text = "[Chart placeholder - Regional sales data to be visualized]"
    np.alignment = PP_ALIGN.CENTER
    nr = np.runs[0]
    nr.font.size = Pt(14)
    nr.font.italic = True
    nr.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    # --- Slide 4: Regional Breakdown ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Regional Breakdown"
    body4 = slide4.placeholders[1].text_frame
    body4.text = "North: $1,520K total ($890K online, $630K in-store)"
    p4a = body4.add_paragraph()
    p4a.text = "South: $1,180K total ($680K online, $500K in-store)"
    p4a.level = 0
    p4b = body4.add_paragraph()
    p4b.text = "East: $1,090K total ($620K online, $470K in-store)"
    p4b.level = 0
    p4c = body4.add_paragraph()
    p4c.text = "West: $1,010K total ($540K online, $470K in-store)"
    p4c.level = 0

    # --- Slide 5: Next Steps ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Next Steps"
    body5 = slide5.placeholders[1].text_frame
    body5.text = "Expand mobile commerce platform in South and West regions"
    p5a = body5.add_paragraph()
    p5a.text = "Launch loyalty program pilot in top-performing North stores"
    p5a.level = 0
    p5b = body5.add_paragraph()
    p5b.text = "Conduct customer survey on channel preferences by Q2 2025"
    p5b.level = 0
    p5c = body5.add_paragraph()
    p5c.text = "Review pricing strategy for in-store vs online alignment"
    p5c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
