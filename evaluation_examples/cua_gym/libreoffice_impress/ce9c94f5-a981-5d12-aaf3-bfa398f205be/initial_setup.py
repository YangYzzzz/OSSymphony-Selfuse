"""
Initial Setup: Create a 9-slide CloudSync Solutions sales pitch presentation.
Task ID: impress_sales_036
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_036'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=RGBColor(0x33, 0x33, 0x33),
                 alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def set_solid_bg(slide, r, g, b):
    """Set solid background color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    DARK_BLUE = RGBColor(0x1A, 0x36, 0x5D)
    MEDIUM_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
    DARK_TEXT = RGBColor(0x2D, 0x3A, 0x4A)
    ACCENT = RGBColor(0x4C, 0x51, 0xBF)

    # ---- Slide 1: Title ----
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_solid_bg(s1, 0x1A, 0x36, 0x5D)
    add_text_box(s1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "CloudSync Solutions", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s1, Inches(2), Inches(3.8), Inches(9), Inches(1.0),
                 "Enterprise Cloud Infrastructure for Modern Teams", font_size=24,
                 color=RGBColor(0xA0, 0xC4, 0xF0), alignment=PP_ALIGN.CENTER)
    add_text_box(s1, Inches(4), Inches(5.5), Inches(5), Inches(0.5),
                 "Q2 2025 Sales Presentation", font_size=16,
                 color=RGBColor(0x80, 0xA0, 0xC0), alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: The Problem ----
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s2, 0xF0, 0xF0, 0xF5)
    add_text_box(s2, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "The Challenge", font_size=36, bold=True, color=DARK_BLUE,
                 alignment=PP_ALIGN.LEFT)
    problems = [
        "87% of enterprises struggle with multi-cloud management complexity",
        "Average downtime costs $5,600 per minute for mid-size companies",
        "Security breaches from misconfigured cloud services rose 43% in 2024",
        "IT teams spend 60% of their time on manual infrastructure tasks",
    ]
    y = Inches(2.0)
    for prob in problems:
        add_text_box(s2, Inches(1.5), y, Inches(10), Inches(0.7),
                     f"  {prob}", font_size=18, color=DARK_TEXT)
        y += Inches(1.0)

    # ---- Slide 3: Our Solution ----
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s3, 0xFF, 0xFF, 0xFF)
    add_text_box(s3, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "CloudSync Platform", font_size=36, bold=True, color=DARK_BLUE)
    add_text_box(s3, Inches(1), Inches(1.8), Inches(10), Inches(1.0),
                 "A unified control plane that simplifies cloud operations across AWS, Azure, and GCP.",
                 font_size=20, color=DARK_TEXT)
    features_brief = [
        "Single dashboard for all cloud resources",
        "AI-powered anomaly detection and auto-remediation",
        "Zero-trust security framework built in",
    ]
    y = Inches(3.2)
    for feat in features_brief:
        add_text_box(s3, Inches(1.5), y, Inches(10), Inches(0.6),
                     f"  {feat}", font_size=18, color=DARK_TEXT)
        y += Inches(0.8)

    # ---- Slide 4: Key Features ----
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s4, 0xF0, 0xF0, 0xF5)
    add_text_box(s4, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "Key Features", font_size=36, bold=True, color=DARK_BLUE)
    features = [
        ("Unified Dashboard", "Monitor and manage 500+ cloud services from one interface"),
        ("Smart Scaling", "ML-driven auto-scaling reduces compute costs by up to 40%"),
        ("Compliance Engine", "Automated SOC2, HIPAA, and GDPR compliance checks"),
        ("Incident Response", "Mean time to resolution reduced from 4 hours to 12 minutes"),
    ]
    y = Inches(2.0)
    for title, desc in features:
        add_text_box(s4, Inches(1.5), y, Inches(4), Inches(0.5),
                     title, font_size=22, bold=True, color=MEDIUM_BLUE)
        add_text_box(s4, Inches(1.5), y + Inches(0.5), Inches(10), Inches(0.5),
                     desc, font_size=16, color=DARK_TEXT)
        y += Inches(1.2)

    # ---- Slide 5: Customer Benefits ----
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s5, 0xFF, 0xFF, 0xFF)
    add_text_box(s5, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "Measurable Impact", font_size=36, bold=True, color=DARK_BLUE)
    metrics = [
        ("40%", "Reduction in Cloud Spend"),
        ("99.99%", "Platform Uptime SLA"),
        ("3x", "Faster Deployment Cycles"),
        ("67%", "Less Manual Configuration"),
    ]
    x = Inches(1.0)
    for val, label in metrics:
        add_text_box(s5, x, Inches(2.5), Inches(2.5), Inches(1.0),
                     val, font_size=48, bold=True, color=ACCENT,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(s5, x, Inches(4.0), Inches(2.5), Inches(0.8),
                     label, font_size=16, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)
        x += Inches(3.0)

    # ---- Slide 6: Case Study ----
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s6, 0x1A, 0x36, 0x5D)
    add_text_box(s6, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "Case Study: NovaTech Industries", font_size=32, bold=True, color=WHITE)
    add_text_box(s6, Inches(1), Inches(2.0), Inches(10), Inches(1.5),
                 '"CloudSync helped us consolidate 12 different cloud tools into one platform. '
                 'We cut our infrastructure costs by $2.3M annually while improving uptime from '
                 '99.5% to 99.99%."',
                 font_size=18, color=RGBColor(0xC0, 0xD0, 0xE8))
    add_text_box(s6, Inches(1), Inches(4.2), Inches(10), Inches(0.5),
                 "- Rachel Torres, VP of Engineering, NovaTech Industries",
                 font_size=16, bold=True, color=RGBColor(0xA0, 0xC4, 0xF0))
    results = [
        "$2.3M saved annually",
        "Deployment time: 3 days to 4 hours",
        "Zero security incidents in 18 months",
    ]
    y = Inches(5.2)
    for r in results:
        add_text_box(s6, Inches(1.5), y, Inches(10), Inches(0.5),
                     f"  {r}", font_size=16, color=WHITE)
        y += Inches(0.6)

    # ---- Slide 7: Pricing ----
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s7, 0xF0, 0xF0, 0xF5)
    add_text_box(s7, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "Pricing Plans", font_size=36, bold=True, color=DARK_BLUE)
    plans = [
        ("Starter", "$2,500/mo", "Up to 50 cloud resources\nBasic monitoring\nEmail support"),
        ("Professional", "$7,500/mo", "Up to 500 resources\nAdvanced analytics\n24/7 support"),
        ("Enterprise", "Custom", "Unlimited resources\nDedicated success manager\nSLA guarantee"),
    ]
    x = Inches(1.0)
    for name, price, details in plans:
        add_text_box(s7, x, Inches(2.0), Inches(3.5), Inches(0.6),
                     name, font_size=24, bold=True, color=DARK_BLUE,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(s7, x, Inches(2.8), Inches(3.5), Inches(0.8),
                     price, font_size=32, bold=True, color=ACCENT,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(s7, x, Inches(4.0), Inches(3.5), Inches(2.0),
                     details, font_size=14, color=DARK_TEXT,
                     alignment=PP_ALIGN.CENTER)
        x += Inches(4.0)

    # ---- Slide 8: Team ----
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s8, 0xFF, 0xFF, 0xFF)
    add_text_box(s8, Inches(1), Inches(0.5), Inches(11), Inches(1.0),
                 "Leadership Team", font_size=36, bold=True, color=DARK_BLUE)
    team = [
        ("David Park", "CEO & Co-Founder", "Former AWS VP, 15 years in cloud infrastructure"),
        ("Priya Sharma", "CTO & Co-Founder", "Ex-Google Cloud architect, PhD in distributed systems"),
        ("James Mitchell", "VP of Sales", "Built $200M ARR sales org at Datadog"),
        ("Lisa Chen", "VP of Engineering", "Former Meta infrastructure lead, 12 years experience"),
    ]
    y = Inches(2.0)
    for name, role, bio in team:
        add_text_box(s8, Inches(1.5), y, Inches(4), Inches(0.5),
                     name, font_size=20, bold=True, color=DARK_BLUE)
        add_text_box(s8, Inches(1.5), y + Inches(0.4), Inches(4), Inches(0.5),
                     role, font_size=16, bold=True, color=ACCENT)
        add_text_box(s8, Inches(5.5), y + Inches(0.1), Inches(7), Inches(0.6),
                     bio, font_size=15, color=DARK_TEXT)
        y += Inches(1.2)

    # ---- Slide 9: Q&A ----
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_solid_bg(s9, 0x1A, 0x36, 0x5D)
    add_text_box(s9, Inches(2), Inches(2.5), Inches(9), Inches(1.5),
                 "Questions & Discussion", font_size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(s9, Inches(3), Inches(4.5), Inches(7), Inches(0.8),
                 "We look forward to answering your questions", font_size=20,
                 color=RGBColor(0xA0, 0xC4, 0xF0), alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready: open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
