"""
Reward Script: Create CTA slide with gradient background and contact info
Task ID: impress_sales_036
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 10 exists as the last slide (0.15)
  Component 2: CTA heading text, 36pt bold white centered (0.30)
  Component 3: Contact info text, 18pt white centered (0.25)
  Component 4: Gradient background left #2B6CB0 to right #4C51BF (0.30)
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_036'

def get_all_text_shapes(slide):
    """Recursively get all text shapes including those in groups."""
    def extract(shape):
        results = []
        if hasattr(shape, "text") and hasattr(shape, "text_frame"):
            results.append(shape)
        if hasattr(shape, 'shapes'):
            for sub in shape.shapes:
                results.extend(extract(sub))
        return results
    out = []
    for shape in slide.shapes:
        out.extend(extract(shape))
    return out


def check_gradient_background(pptx_path, slide_idx):
    """
    Check slide background for gradient with specific colors.
    Returns tuple: (has_gradient, left_color_match, right_color_match)
    """
    ns = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            slide_xml = f'ppt/slides/slide{slide_idx + 1}.xml'
            with zf.open(slide_xml) as f:
                content = f.read()
                root = ET.fromstring(content)

            # bgPr can be nested under p:cSld or at other levels
            # Search broadly for gradFill under bgPr
            bgPr = root.find('.//{%s}bgPr' % ns['p'])
            if bgPr is None:
                print("  DEBUG: No bgPr element found")
                return (False, False, False)

            gradFill = bgPr.find('{%s}gradFill' % ns['a'])
            if gradFill is None:
                print("  DEBUG: No gradFill in bgPr")
                return (False, False, False)

            # Check linear direction (ang=0 means left to right)
            lin = gradFill.find('{%s}lin' % ns['a'])
            if lin is not None:
                ang = lin.get('ang', '')
                print(f"  DEBUG: Gradient linear angle: {ang}")

            gsLst = gradFill.find('{%s}gsLst' % ns['a'])
            if gsLst is None:
                print("  DEBUG: No gsLst found")
                return (True, False, False)

            stops = []
            for gs in gsLst.findall('{%s}gs' % ns['a']):
                pos = gs.get('pos', '')
                srgb = gs.find('{%s}srgbClr' % ns['a'])
                color = srgb.get('val', '').upper() if srgb is not None else ''
                stops.append((pos, color))
                print(f"  DEBUG: Gradient stop pos={pos} color={color}")

            # Check colors: left (#2B6CB0) at pos 0, right (#4C51BF) at pos 100000
            left_ok = False
            right_ok = False
            for pos, color in stops:
                if pos == '0' and color == '2B6CB0':
                    left_ok = True
                if pos == '100000' and color == '4C51BF':
                    right_ok = True

            return (True, left_ok, right_ok)
    except Exception as e:
        print(f"  DEBUG: Gradient check error: {e}")
        return (False, False, False)


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # Component 1: Slide 10 exists as the last slide (0.15 points)
    # Initial has 9 slides, golden has 10. We check that there are at least 10 slides.
    try:
        if num_slides >= 10:
            print(f"PASS: Component 1 -- Slide 10 exists (total slides: {num_slides}) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 1 -- Expected >= 10 slides, found {num_slides}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # If fewer than 10 slides, remaining checks will fail; return early
    if num_slides < 10:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    last_slide = prs.slides[num_slides - 1]
    text_shapes = get_all_text_shapes(last_slide)

    # Component 2: CTA heading "Ready to Transform Your Business?" in 36pt bold white centered (0.30 points)
    try:
        cta_found = False
        cta_score = 0.0

        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                full_text = para.text.strip()
                if "Ready to Transform Your Business?" in full_text:
                    cta_found = True
                    # Sub-check: bold
                    all_bold = True
                    all_white = True
                    size_ok = True
                    for run in para.runs:
                        if not (run.text or "").strip():
                            continue
                        bold_val = run.font.bold
                        if bold_val is None or bold_val is False:
                            all_bold = False
                        # Check white color
                        try:
                            if run.font.color.type is not None:
                                rgb_str = str(run.font.color.rgb).upper()
                                if rgb_str != 'FFFFFF':
                                    all_white = False
                            else:
                                all_white = False
                        except:
                            all_white = False
                        # Check size ~36pt (457200 EMU = 36pt, or sz=3600 in hundredths)
                        if run.font.size is not None:
                            # font.size is in EMU: 36pt = 457200 EMU
                            actual_pt = run.font.size / 12700  # convert EMU to pt
                            if abs(actual_pt - 36) > 2:
                                size_ok = False
                        else:
                            size_ok = False

                    # Check centering
                    from pptx.enum.text import PP_ALIGN
                    centered = (para.alignment == PP_ALIGN.CENTER)

                    if all_bold:
                        cta_score += 0.10
                    else:
                        print(f"  FAIL: CTA text is not bold")
                    if all_white:
                        cta_score += 0.05
                    else:
                        print(f"  FAIL: CTA text is not white")
                    if size_ok:
                        cta_score += 0.10
                    else:
                        print(f"  FAIL: CTA text size is not 36pt")
                    if centered:
                        cta_score += 0.05
                    else:
                        print(f"  FAIL: CTA text is not centered")
                    break
            if cta_found:
                break

        if cta_found:
            total_score += cta_score
            print(f"PASS: Component 2 -- CTA heading found with score {cta_score}/0.30 ({cta_score} pts)")
        else:
            print(f"FAIL: Component 2 -- CTA heading 'Ready to Transform Your Business?' not found on last slide")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Contact info text in 18pt white centered (0.25 points)
    try:
        contact_found = False
        contact_score = 0.0

        for shape in text_shapes:
            for para in shape.text_frame.paragraphs:
                full_text = para.text.strip()
                if "sales@cloudsync.io" in full_text and "(555) 123-4567" in full_text:
                    contact_found = True

                    all_white = True
                    size_ok = True
                    for run in para.runs:
                        if not (run.text or "").strip():
                            continue
                        try:
                            if run.font.color.type is not None:
                                rgb_str = str(run.font.color.rgb).upper()
                                if rgb_str != 'FFFFFF':
                                    all_white = False
                            else:
                                all_white = False
                        except:
                            all_white = False
                        if run.font.size is not None:
                            actual_pt = run.font.size / 12700
                            if abs(actual_pt - 18) > 2:
                                size_ok = False
                        else:
                            size_ok = False

                    from pptx.enum.text import PP_ALIGN
                    centered = (para.alignment == PP_ALIGN.CENTER)

                    if all_white:
                        contact_score += 0.10
                    else:
                        print(f"  FAIL: Contact text is not white")
                    if size_ok:
                        contact_score += 0.10
                    else:
                        print(f"  FAIL: Contact text size is not 18pt")
                    if centered:
                        contact_score += 0.05
                    else:
                        print(f"  FAIL: Contact text is not centered")
                    break
            if contact_found:
                break

        if contact_found:
            total_score += contact_score
            print(f"PASS: Component 3 -- Contact info found with score {contact_score}/0.25 ({contact_score} pts)")
        else:
            print(f"FAIL: Component 3 -- Contact info not found on last slide")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Gradient background left #2B6CB0 to right #4C51BF (0.30 points)
    try:
        has_gradient, left_ok, right_ok = check_gradient_background(file_path, num_slides - 1)

        grad_score = 0.0
        if has_gradient:
            grad_score += 0.10
            print(f"  PASS: Gradient fill present on last slide")
        else:
            print(f"  FAIL: No gradient fill on last slide background")

        if left_ok:
            grad_score += 0.10
            print(f"  PASS: Left gradient color is #2B6CB0")
        else:
            if has_gradient:
                print(f"  FAIL: Left gradient color is not #2B6CB0")

        if right_ok:
            grad_score += 0.10
            print(f"  PASS: Right gradient color is #4C51BF")
        else:
            if has_gradient:
                print(f"  FAIL: Right gradient color is not #4C51BF")

        total_score += grad_score
        print(f"PASS: Component 4 -- Gradient background score {grad_score}/0.30 ({grad_score} pts)")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
