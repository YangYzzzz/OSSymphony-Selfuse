"""
Reward Script: Create comparison table on slide 5
Task ID: impress_stu_021
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Table exists on slide 5 with 5 rows x 3 columns
  Component 2 (0.15): Header row has correct content
  Component 3 (0.15): Row 1 - Functionalism data
  Component 4 (0.15): Row 2 - Marxism data
  Component 5 (0.15): Row 3 - Feminism data
  Component 6 (0.15): Row 4 - Symbolic Interactionism data
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_021'


def persist_app_state():
    """Send Ctrl+S to save any unsaved LibreOffice changes."""
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(1.0)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 5 slides
    if len(prs.slides) < 5:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 5")
        print("REWARD: 0.0")
        return 0.0

    slide5 = prs.slides[4]  # 0-indexed, slide 5

    # Find table on slide 5
    table = None
    for shape in slide5.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            break

    # Component 1: Table exists on slide 5 with correct dimensions (0.25 points)
    try:
        if table is not None:
            nrows = len(table.rows)
            ncols = len(table.columns)
            if nrows == 5 and ncols == 3:
                print(f"PASS: Component 1 - Table found on slide 5 with {nrows}x{ncols} (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - Table dimensions are {nrows}x{ncols}, expected 5x3")
        else:
            print("FAIL: Component 1 - No table found on slide 5")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    if table is None:
        # No table means no further checks can pass
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    # Helper to normalize cell text for comparison
    def norm(text):
        return text.strip().lower() if text else ""

    # Expected data: header + 4 rows
    expected_header = ["Theory", "Key Idea", "Criticism"]
    expected_rows = [
        ["Functionalism", "Society as a system", "Ignores conflict"],
        ["Marxism", "Class struggle", "Too deterministic"],
        ["Feminism", "Gender inequality", "Overlooks intersectionality"],
        ["Symbolic Interactionism", "Micro-level meaning", "Ignores macro structures"],
    ]

    # Component 2: Header row content (0.15 points)
    try:
        actual_header = [table.cell(0, c).text.strip() for c in range(min(len(table.columns), 3))]
        if all(norm(actual_header[i]) == norm(expected_header[i]) for i in range(3)):
            print(f"PASS: Component 2 - Header row matches: {actual_header} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 - Header row {actual_header} != expected {expected_header}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Components 3-6: Data rows (0.15 points each)
    for idx, expected_row in enumerate(expected_rows):
        comp_num = idx + 3
        row_idx = idx + 1  # skip header
        try:
            if row_idx < len(table.rows):
                actual_row = [table.cell(row_idx, c).text.strip() for c in range(min(len(table.columns), 3))]
                if all(norm(actual_row[i]) == norm(expected_row[i]) for i in range(3)):
                    print(f"PASS: Component {comp_num} - Row {row_idx} matches: {actual_row} (0.15 pts)")
                    total_score += 0.15
                else:
                    print(f"FAIL: Component {comp_num} - Row {row_idx} {actual_row} != expected {expected_row}")
            else:
                print(f"FAIL: Component {comp_num} - Row {row_idx} does not exist (table has {len(table.rows)} rows)")
        except Exception as e:
            print(f"ERROR: Component {comp_num} - {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {final_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state()
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
