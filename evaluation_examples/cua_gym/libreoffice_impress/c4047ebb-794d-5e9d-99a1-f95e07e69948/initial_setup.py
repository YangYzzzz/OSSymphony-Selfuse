"""
Initial Setup: Create a sociology theories presentation with 8 slides.
Slide 5 has title 'Comparison of Sociological Theories' but no table.
Task ID: impress_stu_021
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_021'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.text = body_lines[0]
    for line in body_lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    """Add a slide with only a title (layout 5 = blank, add textbox for title)."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    add_title_slide(prs, "Introduction to Sociological Theories",
                    "SOC 201 - Foundations of Sociology\nDr. Elena Rodriguez\nSpring 2025")

    # Slide 2: Course Overview
    add_content_slide(prs, "Course Overview", [
        "Explore major sociological paradigms and their applications",
        "Understand how different theories explain social phenomena",
        "Analyze strengths and limitations of each perspective",
        "Apply theoretical frameworks to contemporary issues",
    ])

    # Slide 3: What is Sociological Theory?
    add_content_slide(prs, "What is Sociological Theory?", [
        "A set of ideas that provides an explanation for human society",
        "Helps us understand patterns in social behavior",
        "Offers frameworks for interpreting social institutions",
        "Bridges the gap between individual experience and social structures",
    ])

    # Slide 4: Major Theoretical Perspectives
    add_content_slide(prs, "Major Theoretical Perspectives", [
        "Functionalism (Structural Functionalism)",
        "Conflict Theory / Marxism",
        "Feminist Theory",
        "Symbolic Interactionism",
        "Each perspective offers unique insights into society",
    ])

    # Slide 5: Comparison - title only, NO TABLE
    add_title_only_slide(prs, "Comparison of Sociological Theories")

    # Slide 6: Applying Theory to Real World
    add_content_slide(prs, "Applying Theory to Real World", [
        "Education system: functionalist vs. conflict perspective",
        "Healthcare disparities through feminist lens",
        "Social media interactions via symbolic interactionism",
        "Economic inequality analyzed through Marxist framework",
    ])

    # Slide 7: Discussion Questions
    add_content_slide(prs, "Discussion Questions", [
        "Which theory best explains income inequality in modern society?",
        "How do feminist and Marxist perspectives overlap?",
        "Can symbolic interactionism address systemic issues?",
        "Is it possible to integrate multiple theoretical frameworks?",
    ])

    # Slide 8: References & Further Reading
    add_content_slide(prs, "References & Further Reading", [
        "Ritzer, G. (2021). Sociological Theory. McGraw-Hill Education.",
        "Giddens, A. (2018). Sociology. Polity Press.",
        "Collins, P.H. (2019). Intersectionality as Critical Social Theory.",
        "Mead, G.H. (1934). Mind, Self, and Society.",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
