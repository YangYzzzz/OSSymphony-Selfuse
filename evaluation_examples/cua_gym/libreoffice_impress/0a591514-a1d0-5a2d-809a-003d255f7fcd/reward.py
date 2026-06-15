"""
Reward Script: Create a title slide with specific title, subtitle, and background color.
Task ID: impress_stu_016
Domain: libreoffice_impress
Scoring:
  - Component 1: Title text matches "Sorting Algorithms Comparison" (0.25)
  - Component 2: Title font is 36pt bold (0.25)
  - Component 3: Subtitle text matches expected (0.25)
  - Component 4: Subtitle font is 16pt (0.10)
  - Component 5: Background color is #F0F0F0 (0.15)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_016'


def persist_app_state(domain):
    """Try to save any unsaved GUI edits via Ctrl+S."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as e:
        print(f"CRITICAL: Missing library: {e}")
        print("REWARD: 0.0")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 1 slide
    if len(prs.slides) == 0:
        print("FAIL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Helper: find shape by name pattern
    title_shape = None
    subtitle_shape = None
    for shape in slide.shapes:
        name_lower = shape.name.lower()
        if 'title' in name_lower and 'subtitle' not in name_lower:
            title_shape = shape
        elif 'subtitle' in name_lower:
            subtitle_shape = shape

    # Fallback: use placeholder indices
    if title_shape is None or subtitle_shape is None:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if title_shape is None:
                    title_shape = shape
                elif subtitle_shape is None:
                    subtitle_shape = shape

    # Component 1: Title text matches (0.25 points)
    try:
        expected_title = "Sorting Algorithms Comparison"
        if title_shape and title_shape.has_text_frame:
            actual_title = title_shape.text_frame.text.strip()
            if actual_title == expected_title:
                print(f"PASS: Component 1 - Title text matches '{actual_title}' (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 1 - Expected title '{expected_title}', found '{actual_title}'")
        else:
            print("FAIL: Component 1 - No title shape found")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Title font is 36pt bold (0.25 points)
    try:
        if title_shape and title_shape.has_text_frame:
            runs = [r for p in title_shape.text_frame.paragraphs for r in p.runs if (r.text or "").strip()]
            if runs:
                # Check bold on any run with text
                any_bold = any(r.font.bold is True for r in runs)
                # Check size - 36pt = 457200 EMU
                any_correct_size = any(r.font.size is not None and r.font.size == 457200 for r in runs)

                bold_pass = any_bold
                size_pass = any_correct_size

                if bold_pass and size_pass:
                    print(f"PASS: Component 2 - Title is 36pt bold (0.25 pts)")
                    total_score += 0.25
                elif bold_pass:
                    print(f"PARTIAL: Component 2 - Title is bold but not 36pt (sizes: {[r.font.size for r in runs]})")
                    total_score += 0.1
                elif size_pass:
                    print(f"PARTIAL: Component 2 - Title is 36pt but not bold (bold: {[r.font.bold for r in runs]})")
                    total_score += 0.1
                else:
                    print(f"FAIL: Component 2 - Title not 36pt bold (size={[r.font.size for r in runs]}, bold={[r.font.bold for r in runs]})")
            else:
                print("FAIL: Component 2 - No text runs in title")
        else:
            print("FAIL: Component 2 - No title shape found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Subtitle text matches (0.25 points)
    try:
        expected_subtitle = "CS 301 - Data Structures | Team: Alex, Priya, Jordan"
        if subtitle_shape and subtitle_shape.has_text_frame:
            actual_subtitle = subtitle_shape.text_frame.text.strip()
            if actual_subtitle == expected_subtitle:
                print(f"PASS: Component 3 - Subtitle text matches (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 3 - Expected subtitle '{expected_subtitle}', found '{actual_subtitle}'")
        else:
            print("FAIL: Component 3 - No subtitle shape found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Subtitle font is 16pt (0.10 points)
    try:
        if subtitle_shape and subtitle_shape.has_text_frame:
            runs = [r for p in subtitle_shape.text_frame.paragraphs for r in p.runs if (r.text or "").strip()]
            if runs:
                # 16pt = 203200 EMU
                any_correct_size = any(r.font.size is not None and r.font.size == 203200 for r in runs)
                if any_correct_size:
                    print(f"PASS: Component 4 - Subtitle is 16pt (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 4 - Subtitle not 16pt (sizes: {[r.font.size for r in runs]})")
            else:
                print("FAIL: Component 4 - No text runs in subtitle")
        else:
            print("FAIL: Component 4 - No subtitle shape found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    # Component 5: Background color is #F0F0F0 (0.15 points)
    try:
        fill = slide.background.fill
        bg_color = None
        if fill.type == 1:  # SOLID
            bg_color = str(fill.fore_color.rgb)
        elif fill.type == 5:  # inherited
            # Check master
            try:
                master_fill = slide.slide_layout.slide_master.background.fill
                if master_fill.type == 1:
                    bg_color = str(master_fill.fore_color.rgb)
            except Exception:
                pass

        if bg_color and bg_color.upper() == "F0F0F0":
            print(f"PASS: Component 5 - Background color is #F0F0F0 (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 - Expected background #F0F0F0, found fill_type={fill.type}, color={bg_color}")
    except Exception as e:
        print(f"ERROR: Component 5 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
