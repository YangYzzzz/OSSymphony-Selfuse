"""
Reward Script: Venture Capital Fund III Presentation
Task ID: impress_wf_085
Domain: libreoffice_impress
Scoring:
  C1: 12 slides (0.15)
  C2: #0A1628 background on all slides (0.10)
  C3: Slide 1 has "Apex Ventures" and "Fund III" (0.10)
  C4: Slide 3 has 4 circle/oval shapes (partner cards) (0.10)
  C5: Slide 4 has a track record table (0.10)
  C6: Slide 5 has 12 rectangle shapes (portfolio grid) (0.10)
  C7: Slide 6 has a chart (portfolio performance) (0.10)
  C8: Slide 8 has 2 charts (market opportunity) (0.10)
  C9: Slide 10 has 5 decreasing-width shapes (funnel) (0.10)
  C10: Gold #C9B037 accents present on multiple slides (0.05)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_wf_085'


def get_slide_bg_color(slide):
    """Get slide background color as uppercase hex string, or None."""
    fill = slide.background.fill
    if fill.type == 1:  # solid fill
        try:
            return str(fill.fore_color.rgb).upper()
        except Exception:
            return None
    elif fill.type == 5:  # inherited from master
        try:
            master_fill = slide.slide_layout.slide_master.background.fill
            if master_fill.type == 1:
                return str(master_fill.fore_color.rgb).upper()
        except Exception:
            pass
    return None


def get_all_text(slide):
    """Get all text from a slide as a single lowercase string."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for r in range(len(table.rows)):
                for c in range(len(table.columns)):
                    texts.append(table.cell(r, c).text)
    return ' '.join(texts).lower()


def has_gold_accent(slide):
    """Check if any shape/text on the slide uses #C9B037 gold color."""
    for shape in slide.shapes:
        # Check auto shape fills
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                if shape.fill.type == 1:
                    if str(shape.fill.fore_color.rgb).upper() == 'C9B037':
                        return True
            except Exception:
                pass
        # Check text colors
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            if str(run.font.color.rgb).upper() == 'C9B037':
                                return True
                    except Exception:
                        pass
        # Check table text colors
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for r in range(len(table.rows)):
                for c in range(len(table.columns)):
                    for para in table.cell(r, c).text_frame.paragraphs:
                        for run in para.runs:
                            try:
                                if run.font.color.type is not None:
                                    if str(run.font.color.rgb).upper() == 'C9B037':
                                        return True
                            except Exception:
                                pass
    return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print("CRITICAL: Cannot load file %s: %s" % (file_path, e))
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)

    # Component 1: Exactly 12 slides (0.15 points)
    try:
        if num_slides == 12:
            print("PASS: Component 1 — 12 slides found (0.15 pts)")
            total_score += 0.15
        else:
            print("FAIL: Component 1 — expected 12 slides, found %d" % num_slides)
    except Exception as e:
        print("ERROR: Component 1 — %s" % e)

    # Need at least 10 slides to check slide-specific components
    if num_slides < 10:
        print("FAIL: Not enough slides for further checks (need >= 10, have %d)" % num_slides)
        final_score = min(total_score, 1.0)
        print("\nScore: %.2f/1.0" % total_score)
        print("REWARD: %.1f" % final_score)
        return final_score

    # Component 2: All slides have #0A1628 background (0.10 points)
    try:
        bg_count = 0
        for slide in slides:
            bg = get_slide_bg_color(slide)
            if bg == '0A1628':
                bg_count += 1
        if bg_count == num_slides:
            print("PASS: Component 2 — all %d slides have #0A1628 background (0.10 pts)" % num_slides)
            total_score += 0.10
        else:
            print("FAIL: Component 2 — %d/%d slides have #0A1628 background" % (bg_count, num_slides))
    except Exception as e:
        print("ERROR: Component 2 — %s" % e)

    # Component 3: Slide 1 contains "Apex Ventures" and "Fund III" (0.10 points)
    try:
        slide1_text = get_all_text(slides[0])
        has_apex = 'apex ventures' in slide1_text
        has_fund = 'fund iii' in slide1_text or 'fund 3' in slide1_text
        if has_apex and has_fund:
            print("PASS: Component 3 — Slide 1 has 'Apex Ventures' and 'Fund III' (0.10 pts)")
            total_score += 0.10
        else:
            print("FAIL: Component 3 — Slide 1 missing title text (apex=%s, fund=%s)" % (has_apex, has_fund))
    except Exception as e:
        print("ERROR: Component 3 — %s" % e)

    # Component 4: Slide 3 has 4 circle/oval shapes (partner photo placeholders) (0.10 points)
    try:
        oval_count = 0
        for shape in slides[2].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                if 'oval' in shape.name.lower() or 'circle' in shape.name.lower() or 'ellipse' in shape.name.lower():
                    oval_count += 1
                # Also check if it's actually an oval by shape XML
                elif hasattr(shape, '_element'):
                    prstGeom = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
                    if prstGeom is not None and prstGeom.get('prst') in ('ellipse', 'oval'):
                        oval_count += 1
        if oval_count >= 4:
            print("PASS: Component 4 — Slide 3 has %d oval/circle shapes (0.10 pts)" % oval_count)
            total_score += 0.10
        else:
            print("FAIL: Component 4 — Slide 3 has %d oval/circle shapes (expected >= 4)" % oval_count)
    except Exception as e:
        print("ERROR: Component 4 — %s" % e)

    # Component 5: Slide 4 has a table (track record) (0.10 points)
    try:
        track_table_rows = 0
        track_table_cols = 0
        for shape in slides[3].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                all_text = ''
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        all_text += table.cell(r, c).text.lower() + ' '
                if ('fund i' in all_text or 'irr' in all_text or 'return' in all_text or 'multiple' in all_text):
                    track_table_rows = len(table.rows)
                    track_table_cols = len(table.columns)
                    break
        if track_table_rows > 0:
            print("PASS: Component 5 — Slide 4 has track record table (%dx%d) (0.10 pts)" % (track_table_rows, track_table_cols))
            total_score += 0.10
        else:
            any_table = any(s.shape_type == MSO_SHAPE_TYPE.TABLE for s in slides[3].shapes)
            if any_table:
                print("FAIL: Component 5 — Slide 4 has table but no track record content")
            else:
                print("FAIL: Component 5 — Slide 4 has no table")
    except Exception as e:
        print("ERROR: Component 5 — %s" % e)

    # Component 6: Slide 5 has 12 rectangle shapes (portfolio grid) (0.10 points)
    try:
        rect_count = 0
        for shape in slides[4].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                name_lower = shape.name.lower()
                if 'rounded rectangle' in name_lower or 'rectangle' in name_lower:
                    # Skip thin separator lines (height < 50000 EMU ~ 0.05 inch)
                    if shape.height > 100000:  # must be a substantial rectangle
                        rect_count += 1
        if rect_count >= 12:
            print("PASS: Component 6 — Slide 5 has %d rectangles (portfolio grid) (0.10 pts)" % rect_count)
            total_score += 0.10
        else:
            print("FAIL: Component 6 — Slide 5 has %d substantial rectangles (expected >= 12)" % rect_count)
    except Exception as e:
        print("ERROR: Component 6 — %s" % e)

    # Component 7: Slide 6 has at least 1 chart (portfolio performance) (0.10 points)
    try:
        chart_count = 0
        for shape in slides[5].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_count += 1
        if chart_count >= 1:
            print("PASS: Component 7 — Slide 6 has %d chart(s) (0.10 pts)" % chart_count)
            total_score += 0.10
        else:
            print("FAIL: Component 7 — Slide 6 has %d charts (expected >= 1)" % chart_count)
    except Exception as e:
        print("ERROR: Component 7 — %s" % e)

    # Component 8: Slide 8 has 2 charts (market opportunity) (0.10 points)
    try:
        if num_slides >= 8:
            chart_count = 0
            for shape in slides[7].shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart_count += 1
            if chart_count >= 2:
                print("PASS: Component 8 — Slide 8 has %d charts (0.10 pts)" % chart_count)
                total_score += 0.10
            else:
                print("FAIL: Component 8 — Slide 8 has %d charts (expected >= 2)" % chart_count)
        else:
            print("FAIL: Component 8 — Not enough slides (need 8, have %d)" % num_slides)
    except Exception as e:
        print("ERROR: Component 8 — %s" % e)

    # Component 9: Slide 10 has 5 decreasing-width shapes (investment funnel) (0.10 points)
    try:
        if num_slides >= 10:
            funnel_shapes = []
            for shape in slides[9].shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    name_lower = shape.name.lower()
                    if 'rounded rectangle' in name_lower or 'rectangle' in name_lower:
                        if shape.height > 100000:  # skip thin separators
                            funnel_shapes.append((shape.top, shape.width))
            # Sort by top position (top to bottom)
            funnel_shapes.sort(key=lambda x: x[0])
            if len(funnel_shapes) >= 5:
                # Check decreasing widths
                widths = [w for _, w in funnel_shapes[:5]]
                is_decreasing = all(widths[i] > widths[i+1] for i in range(4))
                if is_decreasing:
                    print("PASS: Component 9 — Slide 10 has 5 decreasing-width funnel shapes (0.10 pts)")
                    total_score += 0.10
                else:
                    print("FAIL: Component 9 — Slide 10 shapes not in decreasing width order: %s" % widths)
            else:
                print("FAIL: Component 9 — Slide 10 has %d funnel shapes (expected >= 5)" % len(funnel_shapes))
        else:
            print("FAIL: Component 9 — Not enough slides (need 10, have %d)" % num_slides)
    except Exception as e:
        print("ERROR: Component 9 — %s" % e)

    # Component 10: Gold #C9B037 accents present on multiple slides (0.05 points)
    try:
        gold_slide_count = 0
        for slide in slides:
            if has_gold_accent(slide):
                gold_slide_count += 1
        if gold_slide_count >= 3:
            print("PASS: Component 10 — Gold #C9B037 accents found on %d slides (0.05 pts)" % gold_slide_count)
            total_score += 0.05
        else:
            print("FAIL: Component 10 — Gold accents found on only %d slides (expected >= 3)" % gold_slide_count)
    except Exception as e:
        print("ERROR: Component 10 — %s" % e)

    final_score = min(total_score, 1.0)
    print("\nScore: %.2f/1.0" % total_score)
    print("REWARD: %.1f" % final_score)
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print("PERSIST_WARN: save hook failed: %s" % e)


# Main execution
file_path = os.path.join(WORKDIR, 'Fund_III_Deck.pptx')
if not os.path.exists(file_path):
    # Also check Desktop
    alt_path = os.path.join(WORKDIR, 'Desktop', 'Fund_III_Deck.pptx')
    if os.path.exists(alt_path):
        file_path = alt_path
    else:
        print("File not found: %s" % file_path)
        print("REWARD: 0.0")
        exit()

persist_app_state()
verify_task(file_path)
