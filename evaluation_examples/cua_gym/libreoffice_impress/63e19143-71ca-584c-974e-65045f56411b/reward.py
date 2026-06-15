"""
Reward Script: Crescent moon shape via circle subtraction on slide 4
Task ID: impress_ndo_054
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.25): Original two overlapping circles are removed from slide 4
  - Component 2 (0.35): A freeform/custom shape (subtraction result) exists on slide 4
  - Component 3 (0.25): The crescent shape has solid fill color #F1C40F
  - Component 4 (0.15): Shape dimensions are reasonable for a crescent moon
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_054'


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(1.0)
            print("PERSIST: ctrl+s sent for libreoffice_impress")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: file must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Expected at least 4 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)
    shapes = list(slide.shapes)

    # Classify shapes on slide 4
    auto_shape_ovals = []
    freeform_shapes = []
    other_custom_shapes = []

    for shape in shapes:
        # Check for oval/circle AUTO_SHAPEs (the original circles)
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Check if it's an oval by name or by examining the shape
            name_lower = shape.name.lower()
            if 'oval' in name_lower or 'circle' in name_lower or 'ellipse' in name_lower:
                auto_shape_ovals.append(shape)
            else:
                other_custom_shapes.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            freeform_shapes.append(shape)

    print(f"INFO: Slide 4 has {len(shapes)} shapes total")
    print(f"INFO: Found {len(auto_shape_ovals)} oval auto-shapes")
    print(f"INFO: Found {len(freeform_shapes)} freeform shapes")
    print(f"INFO: Found {len(other_custom_shapes)} other auto-shapes")

    # Component 1: Original two overlapping circles are removed (0.25 points)
    # In the initial state, there are 2 ovals. After subtraction, there should be 0.
    try:
        if len(auto_shape_ovals) == 0:
            print(f"PASS: Component 1 -- No oval auto-shapes remain on slide 4 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 -- Found {len(auto_shape_ovals)} oval auto-shapes, expected 0")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: A freeform shape (subtraction result) exists on slide 4 (0.35 points)
    # The Subtract operation in LibreOffice creates a FREEFORM shape
    try:
        if len(freeform_shapes) >= 1:
            print(f"PASS: Component 2 -- Freeform shape found on slide 4 (0.35 pts)")
            total_score += 0.35
        else:
            # Also accept other custom shapes that could be the result of subtraction
            # (in case LibreOffice uses a different shape type)
            has_custom = False
            for shape in other_custom_shapes:
                try:
                    fill = shape.fill
                    if fill.type == 1:  # solid fill
                        rgb = str(fill.fore_color.rgb)
                        if rgb == 'F1C40F':
                            has_custom = True
                            freeform_shapes.append(shape)  # treat as the crescent
                            break
                except:
                    pass
            if has_custom:
                print(f"PASS: Component 2 -- Custom shape with correct fill found on slide 4 (0.35 pts)")
                total_score += 0.35
            else:
                print(f"FAIL: Component 2 -- No freeform/custom shape found on slide 4")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: The crescent shape has solid fill color #F1C40F (0.25 points)
    try:
        crescent_has_correct_fill = False
        for shape in freeform_shapes:
            try:
                fill = shape.fill
                if fill.type == 1:  # solid fill
                    rgb = str(fill.fore_color.rgb)
                    print(f"INFO: Freeform shape '{shape.name}' fill color: {rgb}")
                    if rgb == 'F1C40F':
                        crescent_has_correct_fill = True
                        break
            except Exception as inner_e:
                print(f"INFO: Could not read fill for shape '{shape.name}': {inner_e}")

        if crescent_has_correct_fill:
            print(f"PASS: Component 3 -- Crescent shape fill color is #F1C40F (0.25 pts)")
            total_score += 0.25
        else:
            # Report what colors we found
            for shape in freeform_shapes:
                try:
                    fill = shape.fill
                    if fill.type == 1:
                        print(f"FAIL: Component 3 -- Found fill color {fill.fore_color.rgb}, expected F1C40F")
                    else:
                        print(f"FAIL: Component 3 -- Fill type is {fill.type}, expected solid (1)")
                except:
                    pass
            if len(freeform_shapes) == 0:
                print(f"FAIL: Component 3 -- No freeform shapes to check fill color on")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Shape dimensions are reasonable for a crescent (0.15 points)
    # The crescent should have width < height (since it's a crescent, the subtraction
    # removes part of the width). Golden: w=2160000, h=2880000
    # Initial ovals: large one was 2880000x2880000. Crescent should be narrower.
    try:
        crescent_dims_ok = False
        for shape in freeform_shapes:
            w = shape.width
            h = shape.height
            print(f"INFO: Freeform shape dims: w={w}, h={h}")
            # The crescent width should be less than the original circle diameter (2880000)
            # and the height should be similar to the original circle
            if w > 0 and h > 0 and w < h:
                crescent_dims_ok = True
                break

        if crescent_dims_ok:
            print(f"PASS: Component 4 -- Crescent dimensions are reasonable (width < height) (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 4 -- Crescent dimensions not as expected")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
