"""
Initial Setup: Create Night.pptx with two overlapping circles on slide 4
Task ID: impress_ndo_054
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_054'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Night Sky Photography"
    slide1.placeholders[1].text = "A Guide to Capturing the Stars"
    # Dark background
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)
    # Style title text white
    for run in slide1.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.size = Pt(40)
    for run in slide1.placeholders[1].text_frame.paragraphs[0].runs:
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        run.font.size = Pt(20)

    # --- Slide 2: Photography Tips ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill2 = slide2.background.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txBox2 = slide2.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Essential Night Photography Tips"
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.runs[0]
    run2.font.name = "Arial"
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    tips = [
        "Use a sturdy tripod to avoid camera shake during long exposures",
        "Set ISO between 1600-3200 for optimal star visibility",
        "Use a wide-angle lens with aperture f/2.8 or wider",
        "Enable mirror lock-up to reduce vibration",
        "Shoot in RAW format for maximum post-processing flexibility",
    ]
    txBox2b = slide2.shapes.add_textbox(Cm(2), Cm(4), Cm(20), Cm(12))
    tf2b = txBox2b.text_frame
    tf2b.word_wrap = True
    for i, tip in enumerate(tips):
        p = tf2b.paragraphs[0] if i == 0 else tf2b.add_paragraph()
        p.text = tip
        p.space_after = Pt(12)
        for r in p.runs:
            r.font.name = "Arial"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # --- Slide 3: Equipment ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    txBox3 = slide3.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Recommended Equipment"
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.runs[0]
    run3.font.name = "Arial"
    run3.font.size = Pt(28)
    run3.font.bold = True
    run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add table with equipment
    rows, cols = 5, 3
    tbl_shape = slide3.shapes.add_table(rows, cols, Cm(2), Cm(4), Cm(20), Cm(8))
    tbl = tbl_shape.table
    headers = ["Equipment", "Model", "Price Range"]
    equipment_data = [
        ["Camera Body", "Canon EOS R5", "$3,200 - $3,900"],
        ["Wide Lens", "Sigma 14mm f/1.8 Art", "$1,400 - $1,600"],
        ["Tripod", "Gitzo Systematic GT3543LS", "$850 - $1,000"],
        ["Star Tracker", "Sky-Watcher Star Adventurer", "$350 - $450"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(14)
    for r, row_data in enumerate(equipment_data, 1):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                run.font.size = Pt(12)

    # --- Slide 4: Two Overlapping Circles (the task slide) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill4 = slide4.background.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)  # Dark navy

    # Title for slide 4
    txBox4 = slide4.shapes.add_textbox(Cm(2), Cm(0.5), Cm(20), Cm(2))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "Moon Shape Construction"
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.runs[0]
    run4.font.name = "Arial"
    run4.font.size = Pt(24)
    run4.font.bold = True
    run4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Larger white circle: 8cm diameter, centered roughly on slide
    # Position: centered horizontally, centered vertically
    large_d = Cm(8)
    large_left = Cm(8)    # position on slide
    large_top = Cm(5)
    circle_large = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, large_left, large_top, large_d, large_d
    )
    circle_large.fill.solid()
    circle_large.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    circle_large.line.fill.background()  # no outline

    # Smaller gray circle: 7cm diameter, offset right by 2cm
    small_d = Cm(7)
    small_left = large_left + Cm(2)  # offset 2cm to the right
    # Vertically center the smaller circle relative to the larger one
    small_top = large_top + (large_d - small_d) // 2
    circle_small = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, small_left, small_top, small_d, small_d
    )
    circle_small.fill.solid()
    circle_small.fill.fore_color.rgb = RGBColor(0x99, 0x99, 0x99)
    circle_small.line.fill.background()  # no outline

    # --- Slide 5: Summary ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    fill5 = slide5.background.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    txBox5 = slide5.shapes.add_textbox(Cm(3), Cm(3), Cm(18), Cm(4))
    tf5 = txBox5.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "Happy Shooting!"
    p5.alignment = PP_ALIGN.CENTER
    run5 = p5.runs[0]
    run5.font.name = "Arial"
    run5.font.size = Pt(36)
    run5.font.bold = True
    run5.font.color.rgb = RGBColor(0xF1, 0xC4, 0x0F)

    p5b = tf5.add_paragraph()
    p5b.text = "Remember: patience and practice make perfect night sky photographs."
    p5b.alignment = PP_ALIGN.CENTER
    for r in p5b.runs:
        r.font.name = "Arial"
        r.font.size = Pt(18)
        r.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # Launch LibreOffice Impress with the file
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
