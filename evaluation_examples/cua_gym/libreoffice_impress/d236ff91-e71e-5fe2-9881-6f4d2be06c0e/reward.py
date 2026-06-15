"""
FINAL REWARD SCRIPT - SUCCESS
Task: I'm rehearsing for a live demo and need to remind myself that the demo video has to be done playing before I can talk. How can I jot this down as a note in bold on slide 6 in LibreOffice Impress?
Generated: 2025-08-07 09:54:26
Status: success
Model: o4-mini
Total Steps: 2
"""

import os
from pptx import Presentation

def verify_impress_note_bold(file_path):
    """
    Verifies that slide 6 of the given PPTX file contains a bold note with the reminder phrase.
    Scoring:
      - 0.2: File exists
      - 0.2: At least 6 slides present
      - 0.2: Notes slide exists for slide 6
      - 0.2: Target phrase found in notes
      - 0.2: Note text formatted in bold
    Returns a score between 0.0 and 1.0.
    """
    print(f"Checking file existence: {file_path}")
    score = 0.0
    # Requirement 1: File existence
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        print(f"REWARD: {score}")
        return score
    print("✓ File exists (0.2)")
    score += 0.2

    try:
        # Load presentation
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        print(f"Loaded presentation, slide count: {slide_count}")

        # Requirement 2: At least 6 slides
        if slide_count >= 6:
            print("✓ Slide count >= 6 (0.2)")
            score += 0.2
            slide6 = prs.slides[5]
        else:
            print(f"✗ Slide count < 6: {slide_count}")
            print(f"REWARD: {min(score, 1.0)}")
            return min(score, 1.0)

        # Requirement 3: Notes slide exists
        if not hasattr(slide6, 'notes_slide') or slide6.notes_slide is None:
            print("✗ No notes on slide 6")
            print(f"REWARD: {min(score, 1.0)}")
            return min(score, 1.0)
        print("✓ Notes slide exists (0.2)")
        score += 0.2

        # Requirement 4 & 5: Target phrase and bold formatting
        found_text = False
        found_bold = False
        target_phrase = "demo video"

        for shape in slide6.notes_slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                full_text = ''.join(run.text for run in para.runs)
                # Case-insensitive search for phrase
                if target_phrase in full_text.lower():
                    found_text = True
                    print(f"✓ Found target phrase in notes: '{full_text}' (0.2)")
                    score += 0.2
                    # Check bold formatting
                    for run in para.runs:
                        if run.text and hasattr(run, 'font') and run.font.bold:
                            found_bold = True
                            break
                    if found_bold:
                        print("✓ Found bold formatting for the note (0.2)")
                        score += 0.2
                    else:
                        print("✗ Note text not bold")
                    break
            if found_text:
                break

        if not found_text:
            print("✗ Target phrase not found in notes")

        final_score = min(score, 1.0)
        print(f"REWARD: {final_score}")
        return final_score

    except Exception as e:
        print(f"✗ Error reading presentation: {e}")
        print(f"REWARD: {min(score, 1.0)}")
        return min(score, 1.0)

# Execute verification on the target file
if __name__ == '__main__':
    file_path = "/home/user/im_rehearsing_for_a_live_demo_and_need_to_remind_myself_that_the_demo_video_has_to_be_done_playing_b.pptx"
    verify_impress_note_bold(file_path)

