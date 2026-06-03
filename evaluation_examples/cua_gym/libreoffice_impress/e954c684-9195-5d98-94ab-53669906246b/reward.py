"""
Reward Script: Change slide 2 and slide 6 title text and colors
Task ID: osworld_impress_title_color_match_004
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 title text == 'Background and Context'       (0.3 pts)
  Component 2: Slide 2 title color == #FF6600 (orange, matching slide 4) (0.2 pts)
  Component 3: Slide 6 title text == 'Future Outlook'               (0.3 pts)
  Component 4: Slide 6 title color == #228B22 (green, matching slide 1)  (0.2 pts)
  Total: 1.0
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_004'

# Expected values from task description and context
SLIDE2_EXPECTED_TEXT = 'Background and Context'
SLIDE2_EXPECTED_COLOR = 'FF6600'  # Orange — matching slide 4's title color

SLIDE6_EXPECTED_TEXT = 'Future Outlook'
SLIDE6_EXPECTED_COLOR = '228B22'  # Green — matching slide 1's title color


def get_title_shape(slide):
    """Return the title placeholder shape for the given slide, or None."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0 and shape.has_text_frame:
                return shape
    return None


def get_title_text(slide):
    """Return the title text for a slide, or None if no title found."""
    shape = get_title_shape(slide)
    if shape is None:
        return None
    return shape.text_frame.text


def get_title_colors(slide):
    """
    Return a list of color hex strings from all non-empty runs in the title placeholder.
    Returns empty list if no title or no colored runs.
    """
    shape = get_title_shape(slide)
    if shape is None:
        return []
    colors = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not (run.text or '').strip():
                continue
            try:
                if run.font.color.type is not None:
                    colors.append(str(run.font.color.rgb).upper())
            except Exception:
                pass
    return colors


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"CRITICAL: Expected at least 7 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]  # 0-indexed: slide 2
    slide6 = prs.slides[5]  # 0-indexed: slide 6

    # Component 1: Slide 2 title text is 'Background and Context' (0.3 points)
    try:
        actual_text = get_title_text(slide2)
        if actual_text == SLIDE2_EXPECTED_TEXT:
            print(f"PASS: Component 1 — Slide 2 title text is '{SLIDE2_EXPECTED_TEXT}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 1 — Slide 2 title text expected '{SLIDE2_EXPECTED_TEXT}', found '{actual_text}'")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 2 title color is #FF6600 (orange, matching slide 4) (0.2 points)
    try:
        colors = get_title_colors(slide2)
        if SLIDE2_EXPECTED_COLOR in colors:
            print(f"PASS: Component 2 — Slide 2 title color is #{SLIDE2_EXPECTED_COLOR} (orange) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 — Slide 2 title color expected #{SLIDE2_EXPECTED_COLOR}, found {colors}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 6 title text is 'Future Outlook' (0.3 points)
    try:
        actual_text = get_title_text(slide6)
        if actual_text == SLIDE6_EXPECTED_TEXT:
            print(f"PASS: Component 3 — Slide 6 title text is '{SLIDE6_EXPECTED_TEXT}' (0.3 pts)")
            total_score += 0.3
        else:
            print(f"FAIL: Component 3 — Slide 6 title text expected '{SLIDE6_EXPECTED_TEXT}', found '{actual_text}'")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Slide 6 title color is #228B22 (green, matching slide 1) (0.2 points)
    try:
        colors = get_title_colors(slide6)
        if SLIDE6_EXPECTED_COLOR in colors:
            print(f"PASS: Component 4 — Slide 6 title color is #{SLIDE6_EXPECTED_COLOR} (green) (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 4 — Slide 6 title color expected #{SLIDE6_EXPECTED_COLOR}, found {colors}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
