"""
Initial Setup: 7-slide strategic analysis presentation with specific title colors
Task ID: osworld_impress_title_color_match_004
Domain: libreoffice_impress

Creates a 7-slide strategic analysis deck where:
- Slide 1 title is in green (#228B22)
- Slide 2 title is "Background" in black
- Slide 4 title is in orange (#FF6600)
- Slide 6 title is "Outlook" in black
- Other slides have regular black titles
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_title_color_match_004'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def set_title_color(slide, color_rgb):
    """Set the title placeholder's text color for a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.shape_type == 13:
            continue
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:  # Title placeholder
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                return shape
    return None


def set_title_text_and_color(slide, text, color_rgb):
    """Set the title placeholder's text and color for a slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:  # Title placeholder
                shape.text_frame.paragraphs[0].runs[0].text = text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color_rgb
                return
    # If no runs yet, use title property
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                tf = shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = text
                run.font.color.rgb = color_rgb
                run.font.size = Pt(28)
                return


def add_slide_with_title_content(prs, layout_idx, title_text, title_color, body_text):
    """Add a slide with a title and content."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Set title
    title_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 0:
                title_shape = shape
                break

    if title_shape is not None:
        tf = title_shape.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = title_text
        run.font.color.rgb = title_color
        run.font.size = Pt(32)
        run.font.bold = True

    # Set content body
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            if shape.placeholder_format.idx == 1:  # Content placeholder
                shape.text_frame.text = body_text
                break

    return slide


def create_initial():
    prs = Presentation()
    # Use standard 10x7.5 inch widescreen
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Color definitions
    GREEN = RGBColor(0x22, 0x8B, 0x22)    # #228B22 Forest Green
    ORANGE = RGBColor(0xFF, 0x66, 0x00)   # #FF6600 Orange
    BLACK = RGBColor(0x00, 0x00, 0x00)    # #000000 Black

    # Slide data: (title_text, title_color, body_text)
    slide_data = [
        # Slide 1 - Green title
        (
            "Strategic Analysis 2025",
            GREEN,
            "Executive Summary\nGlobal Market Positioning and Growth Strategy"
        ),
        # Slide 2 - Black title "Background" (to be changed by agent)
        (
            "Background",
            BLACK,
            "Historical Context\n"
            "• Founded in 2010 with focus on enterprise solutions\n"
            "• Expanded to 15 markets by 2018\n"
            "• Revenue grew from $2.3M to $148M over 12 years\n"
            "• Strategic pivot to SaaS model in 2021"
        ),
        # Slide 3 - Black title
        (
            "Market Analysis",
            BLACK,
            "Current Market Landscape\n"
            "• Total addressable market: $12.4B\n"
            "• Serviceable addressable market: $3.8B\n"
            "• Year-over-year growth rate: 18.3%\n"
            "• Key competitor share: 34% combined"
        ),
        # Slide 4 - Orange title
        (
            "Key Challenges",
            ORANGE,
            "Strategic Obstacles\n"
            "• Supply chain disruptions impacting Q3-Q4 delivery\n"
            "• Talent acquisition in engineering: 23 open roles\n"
            "• Regulatory changes in EU market (GDPR v2)\n"
            "• Price pressure from emerging low-cost competitors"
        ),
        # Slide 5 - Black title
        (
            "Opportunities",
            BLACK,
            "Growth Vectors\n"
            "• Asia Pacific expansion: projected $45M revenue\n"
            "• AI/ML product integration roadmap\n"
            "• Strategic acquisitions pipeline: 3 targets identified\n"
            "• New enterprise verticals: healthcare and finance"
        ),
        # Slide 6 - Black title "Outlook" (to be changed by agent)
        (
            "Outlook",
            BLACK,
            "Forward-Looking Projections\n"
            "• FY2026 revenue target: $210M (+42%)\n"
            "• Market share goal: 18% by Q4 2026\n"
            "• New product launches: 4 planned\n"
            "• Geographic expansion: 5 new countries"
        ),
        # Slide 7 - Black title
        (
            "Recommendations",
            BLACK,
            "Action Plan\n"
            "• Accelerate digital transformation initiatives\n"
            "• Invest $8M in R&D for AI capabilities\n"
            "• Restructure go-to-market for SMB segment\n"
            "• Establish Center of Excellence in Singapore"
        ),
    ]

    # Layout 1 = Title and Content layout
    for title_text, title_color, body_text in slide_data:
        add_slide_with_title_content(prs, 1, title_text, title_color, body_text)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
