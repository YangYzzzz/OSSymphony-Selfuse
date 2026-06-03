"""
FINAL REWARD SCRIPT - SUCCESS
Task: On slide 2, I want the Title text box to sit flush against the bottom edge while staying exactly centered left-to-right. How do I reposition it like that in LibreOffice Impress?
Generated: 2025-09-10 13:39:16
Status: success
Model: azure-o3
Total Steps: 2
"""

import os
from pptx import Presentation


def verify_title_position(file_path: str, slide_number: int = 2, tolerance_ratio: float = 0.01) -> float:
    """Verify that on the specified slide (default = 2) the title text box
    • sits flush against the bottom edge, and
    • is horizontally centred.

    A small tolerance (default 1 % of slide dimension) is allowed.

    Returns a progressive score between 0.0 and 1.0.
    """

    max_score = 1.0
    score = 0.0

    # ---------- 1. Load presentation ----------
    if not os.path.exists(file_path):
        print(f"✗ File not found: {file_path}")
        return 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"✗ Could not load presentation: {e}")
        return 0.0

    # ---------- 2. Locate requested slide ----------
    slide_index = slide_number - 1  # convert to 0-based index
    if len(prs.slides) <= slide_index:
        print(f"✗ Presentation only has {len(prs.slides)} slides; slide {slide_number} missing")
        return 0.0

    slide = prs.slides[slide_index]

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    tol_horiz = slide_w * tolerance_ratio   # horizontal tolerance (in EMUs)
    tol_vert = slide_h * tolerance_ratio    # vertical tolerance (in EMUs)

    bottom_ok = False
    centre_ok = False

    # ---------- 3. Inspect text-containing shapes ----------
    for shape in slide.shapes:
        # Only consider shapes that actually contain text
        if not getattr(shape, "has_text_frame", False):
            continue
        if not shape.text or not shape.text.strip():  # ignore empty text boxes
            continue

        text_preview = shape.text.strip()[:30]

        # Horizontal: distance between shape.left and ideal centred left position
        ideal_left = (slide_w - shape.width) / 2
        horiz_gap = abs(shape.left - ideal_left)
        # Vertical: distance between shape bottom and slide bottom
        vert_gap = abs((shape.top + shape.height) - slide_h)

        print(f"Examining shape '{text_preview}' | horiz_gap={horiz_gap} | vert_gap={vert_gap}")

        if vert_gap <= tol_vert:
            bottom_ok = True
            print(f"✓ '{text_preview}' flush with bottom (gap {vert_gap})")
        if horiz_gap <= tol_horiz:
            centre_ok = True
            print(f"✓ '{text_preview}' horizontally centred (gap {horiz_gap})")

        # If both conditions met for any shape, we can stop early
        if bottom_ok and centre_ok:
            break

    # ---------- 4. Scoring ----------
    if bottom_ok:
        score += 0.5
    else:
        print("✗ No text box flush against bottom edge within tolerance")

    if centre_ok:
        score += 0.5
    else:
        print("✗ No text box horizontally centred within tolerance")

    score = min(score, max_score)
    print(f"Total score: {score}/{max_score}")
    return score


if __name__ == "__main__":
    FILE_PATH = "/home/user/on_slide_2_i_want_the_title_text_box_to_sit_flush_against_the_bottom_edge_while_staying_exactly_cent_golden.pptx"
    reward_score = verify_title_position(FILE_PATH)
    print(f"REWARD: {reward_score}")

