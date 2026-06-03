"""
FINAL REWARD SCRIPT - SUCCESS
Task: I’ve got the outline view set up in LibreOffice Impress, and I want it to automatically create a “Summary” slide from that outline. It should land immediately after my Agenda slide (currently slide 2). How do I make Impress do that for me?
Generated: 2025-09-10 14:33:24
Status: success
Model: azure-o3
Total Steps: 4
"""

import os
from pptx import Presentation

def get_slide_title(slide):
    """Return the first non-empty text line on a slide (used as its title)."""
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            text = shape.text_frame.text
            if text:
                for line in text.split("\n"):
                    if line.strip():
                        return line.strip()
    return ""

def verify_impress_summary(file_path):
    """Reward script to verify automatic “Summary” slide generation in LibreOffice Impress.

    Scoring (progressive – max 1.0):
      • 0.4  Summary slide exists immediately after Agenda slide
      • 0.3  Summary slide contains at least two bullet items
      • 0.3  Bullet items match the titles of later slides (awarded proportionally)
    """
    max_score = 1.0
    score = 0.0

    # 1. Load presentation ----------------------------------------------------
    if not (os.path.exists(file_path) and file_path.lower().endswith(".pptx")):
        print(f"✗ Presentation not found or wrong format: {file_path}")
        print(f"REWARD: {score}")
        return score

    try:
        prs = Presentation(file_path)
        print(f"✓ Loaded PPTX with {len(prs.slides)} slides")
    except Exception as e:
        print(f"✗ Failed to load presentation: {e}")
        print(f"REWARD: {score}")
        return score

    # Pre-compute all slide titles
    slide_titles = [get_slide_title(slide) for slide in prs.slides]

    # 2. Locate Agenda slide ---------------------------------------------------
    agenda_idx = None
    for idx, title in enumerate(slide_titles):
        if title and "agenda" in title.lower():
            agenda_idx = idx
            break

    if agenda_idx is None:
        print("✗ No Agenda slide found – cannot verify Summary placement")
    else:
        print(f"✓ Agenda slide found at position {agenda_idx + 1}")

        # 3. Verify Summary slide placement -----------------------------------
        summary_idx = agenda_idx + 1 if agenda_idx + 1 < len(prs.slides) else None
        if summary_idx is not None:
            if slide_titles[summary_idx] and "summary" in slide_titles[summary_idx].lower():
                print(f"✓ Summary slide correctly located immediately after Agenda (slide {summary_idx + 1})")
                score += 0.4
            else:
                print("✗ Slide after Agenda is not a Summary slide")
        else:
            print("✗ Agenda is the last slide; no slide follows it")

    # 4. Inspect Summary slide content ---------------------------------------
    summary_slide = None
    if agenda_idx is not None and agenda_idx + 1 < len(prs.slides):
        candidate = prs.slides[agenda_idx + 1]
        if "summary" in get_slide_title(candidate).lower():
            summary_slide = candidate

    bullet_items = []
    if summary_slide is not None:
        for shape in summary_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text or ""
                lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
                # Remove the title line itself if it is "Summary"
                lines = [ln for ln in lines if ln.lower() != "summary"]
                bullet_items.extend(lines)
        # De-duplicate while preserving order
        seen = set()
        bullet_items = [x for x in bullet_items if not (x in seen or seen.add(x))]

        if len(bullet_items) >= 2:
            print(f"✓ Summary slide contains bullet items: {bullet_items}")
            score += 0.3
        elif bullet_items:
            print("✗ Summary slide has fewer than 2 bullet items – partial credit")
            score += 0.1
        else:
            print("✗ No bullet items found on Summary slide")
    else:
        print("✗ Summary slide not available for content checks")

    # 5. Compare bullet items to later slide titles ---------------------------
    if bullet_items:
        later_titles = [title for i, title in enumerate(slide_titles) if i > (agenda_idx + 1 if agenda_idx is not None else -1) and title]
        match_count = 0
        for item in bullet_items:
            for lt in later_titles:
                if item.lower() == lt.lower():
                    match_count += 1
                    break
        if later_titles:
            match_ratio = match_count / len(later_titles)
            extra_score = 0.3 * match_ratio  # proportional credit
            if match_count:
                print(f"✓ {match_count} of {len(later_titles)} later slide titles are referenced in bullet list (score +{extra_score:.2f})")
            else:
                print("✗ Bullet list does not reference any later slide titles")
            score += extra_score

    # 6. Finalise -------------------------------------------------------------
    final_score = round(min(score, max_score), 2)
    print(f"Total Score: {final_score}")
    print(f"REWARD: {final_score}")
    return final_score

# ---------------- Script entry point ----------------
if __name__ == "__main__":
    PATH = "/home/user/ive_got_the_outline_view_set_up_in_libreoffice_impress_and_i_want_it_to_automatically_create_a_summa_golden.pptx"
    verify_impress_summary(PATH)
