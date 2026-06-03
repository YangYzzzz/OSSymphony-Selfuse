"""
Initial Setup: Create a 7-slide Survey Results presentation with empty content on slide 4.
Task ID: impress_stu_010
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_010'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Survey Results"
    slide1.placeholders[1].text = "Student Engagement & Online Learning Preferences\nSpring 2025"

    # --- Slide 2: Introduction ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Introduction"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Background"
    p = tf2.add_paragraph()
    p.text = "Online learning platforms have become increasingly prevalent in higher education since 2020."
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Research Objective"
    p = tf2.add_paragraph()
    p.text = "This study investigates student preferences for synchronous vs. asynchronous learning modalities."
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Significance"
    p = tf2.add_paragraph()
    p.text = "Findings will help universities optimize course delivery formats for improved student outcomes."
    p.level = 1

    # --- Slide 3: Demographics ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Demographics"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Gender distribution: 58% female, 39% male, 3% non-binary"
    p = tf3.add_paragraph()
    p.text = "Year of study: 22% freshmen, 31% sophomores, 28% juniors, 19% seniors"
    p = tf3.add_paragraph()
    p.text = "Major categories: STEM (45%), Humanities (30%), Business (25%)"
    p = tf3.add_paragraph()
    p.text = "Average GPA: 3.2 (SD = 0.6)"
    p = tf3.add_paragraph()
    p.text = "Prior online course experience: 78% had taken at least one online course"

    # --- Slide 4: Methodology (EMPTY content placeholder) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Methodology"
    # Leave the content placeholder EMPTY - this is where the agent must add the bulleted list
    # Clear any default text
    tf4 = slide4.placeholders[1].text_frame
    tf4.clear()

    # --- Slide 5: Key Findings ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Findings"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "67% of students prefer a hybrid learning model combining online and in-person sessions"
    p = tf5.add_paragraph()
    p.text = "Asynchronous lecture recordings were rated most valuable (4.3/5.0)"
    p = tf5.add_paragraph()
    p.text = "Students in STEM disciplines showed stronger preference for in-person lab components"
    p = tf5.add_paragraph()
    p.text = "Discussion forums were the least preferred online tool (2.8/5.0)"
    p = tf5.add_paragraph()
    p.text = "Internet connectivity was cited as the top barrier to online participation (34%)"

    # --- Slide 6: Discussion ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Discussion"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "The strong preference for hybrid models aligns with recent meta-analyses (Chen et al., 2024)"
    p = tf6.add_paragraph()
    p.text = "Asynchronous content allows students to learn at their own pace, supporting diverse schedules"
    p = tf6.add_paragraph()
    p.text = "STEM preference for in-person labs suggests hands-on components remain essential"
    p = tf6.add_paragraph()
    p.text = "Low engagement with discussion forums may reflect design issues rather than format limitations"

    # --- Slide 7: Conclusion ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "Conclusion"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Universities should invest in hybrid course infrastructure to meet student preferences"
    p = tf7.add_paragraph()
    p.text = "Recorded lectures and flexible scheduling are key drivers of student satisfaction"
    p = tf7.add_paragraph()
    p.text = "Future research should explore long-term academic outcomes across modalities"
    p = tf7.add_paragraph()
    p.text = "Limitations: Single institution, self-reported data, limited to undergraduate population"

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
