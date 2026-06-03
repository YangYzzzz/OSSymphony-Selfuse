"""
Reward Script: Add bulleted list to slide 4 with 5 methodology key points
Task ID: impress_stu_010
Domain: libreoffice_impress
Scoring:
  Component 1: Content placeholder on slide 4 is non-empty (0.2 pts)
  Component 2: Exactly 5 non-empty bullet items (0.2 pts)
  Component 3-7: Each of the 5 required texts present (0.12 pts each, total 0.6 pts)
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_stu_010'

REQUIRED_BULLETS = [
    'Sample size: 150 participants',
    'Age range: 18-25 years',
    'Method: Online survey via Qualtrics',
    'Duration: 2 weeks',
    'Response rate: 72%',
]


def normalize(text):
    """Normalize text for comparison: strip whitespace, collapse internal spaces."""
    return ' '.join(text.strip().split())


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide4 = prs.slides[3]  # 0-indexed

    # Find the content placeholder (not the title)
    content_shape = None
    for shape in slide4.shapes:
        if shape.has_text_frame and shape.name != 'Title 1':
            # Skip the title shape; take the content placeholder
            if 'title' not in shape.name.lower():
                content_shape = shape
                break

    if content_shape is None:
        print("FAIL: No content placeholder found on slide 4")
        print("REWARD: 0.0")
        return 0.0

    # Extract non-empty paragraphs from the content placeholder
    paragraphs = []
    for para in content_shape.text_frame.paragraphs:
        text = normalize(para.text)
        if text:
            paragraphs.append(text)

    print(f"INFO: Found {len(paragraphs)} non-empty paragraphs on slide 4 content placeholder")
    for i, p in enumerate(paragraphs):
        print(f"  [{i}] {repr(p)}")

    # Component 1: Content placeholder has text (not empty) (0.2 points)
    # This differentiates golden (has text) from initial (empty placeholder)
    try:
        if len(paragraphs) > 0:
            print(f"PASS: Component 1 -- Content placeholder has text ({len(paragraphs)} paragraphs) (0.2 pts)")
            total_score += 0.2
        else:
            print("FAIL: Component 1 -- Content placeholder is empty")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Exactly 5 non-empty bullet items (0.2 points)
    try:
        if len(paragraphs) == 5:
            print(f"PASS: Component 2 -- Exactly 5 bullet items found (0.2 pts)")
            total_score += 0.2
        else:
            print(f"FAIL: Component 2 -- Expected 5 bullet items, found {len(paragraphs)}")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Components 3-7: Each required bullet text is present (0.12 points each)
    for idx, expected in enumerate(REQUIRED_BULLETS):
        comp_num = idx + 3
        try:
            expected_norm = normalize(expected)
            found = False
            for para_text in paragraphs:
                if expected_norm.lower() == para_text.lower():
                    found = True
                    break
            if found:
                print(f"PASS: Component {comp_num} -- Found '{expected}' (0.12 pts)")
                total_score += 0.12
            else:
                print(f"FAIL: Component {comp_num} -- '{expected}' not found in slide 4 content")
        except Exception as e:
            print(f"ERROR: Component {comp_num} -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
