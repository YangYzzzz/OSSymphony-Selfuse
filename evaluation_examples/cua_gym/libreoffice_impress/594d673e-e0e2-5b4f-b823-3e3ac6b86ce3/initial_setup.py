"""
Initial Setup: Create a 9-slide sales pitch presentation for CloudSync Pro
Task ID: impress_sales_078
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_078'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Helper to add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
    ACCENT_BLUE = RGBColor(0x2B, 0x6C, 0xB0)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
    BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)

    # ---- Slide 1: Title Slide ----
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = slide1.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    add_text_box(slide1, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "CloudSync Pro", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(3.5), Inches(10), Inches(1.0),
                 "Enterprise Cloud Infrastructure Platform", font_size=24, color=RGBColor(0xBB, 0xD5, 0xED),
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide1, Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                 "Q2 2025 Sales Presentation | Prepared for Meridian Corp", font_size=14,
                 color=RGBColor(0x88, 0xAA, 0xCC), alignment=PP_ALIGN.CENTER)

    # ---- Slide 2: The Challenge ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "The Challenge", font_size=32, bold=True, color=DARK_BLUE)
    challenges = [
        "78% of enterprises report data silos across 5+ cloud providers",
        "Average $2.4M annual spend on redundant cloud management tools",
        "IT teams spend 35% of time on manual sync and migration tasks",
        "Security compliance gaps in multi-cloud environments cost $4.1M per breach",
    ]
    y = Inches(1.8)
    for ch in challenges:
        add_text_box(slide2, Inches(1.2), y, Inches(10), Inches(0.6),
                     f"• {ch}", font_size=16, color=DARK_GRAY)
        y += Inches(0.9)

    # ---- Slide 3: Our Solution ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    fill3 = slide3.background.fill
    fill3.solid()
    fill3.fore_color.rgb = BG_LIGHT
    add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "Our Solution: CloudSync Pro", font_size=32, bold=True, color=DARK_BLUE)
    add_text_box(slide3, Inches(0.8), Inches(1.6), Inches(11), Inches(1.0),
                 "A unified platform that connects, synchronizes, and secures your entire cloud "
                 "ecosystem — reducing operational overhead by 60% while maintaining SOC 2 and "
                 "ISO 27001 compliance.", font_size=16, color=DARK_GRAY)
    features_brief = [
        ("Unified Dashboard", "Single pane of glass for AWS, Azure, GCP, and 20+ cloud services"),
        ("Auto-Sync Engine", "Real-time data synchronization with conflict resolution"),
        ("Smart Migration", "AI-powered workload placement and migration planning"),
    ]
    y = Inches(3.2)
    for title, desc in features_brief:
        add_text_box(slide3, Inches(1.2), y, Inches(4), Inches(0.5),
                     title, font_size=18, bold=True, color=ACCENT_BLUE)
        add_text_box(slide3, Inches(1.2), y + Inches(0.5), Inches(10), Inches(0.5),
                     desc, font_size=14, color=LIGHT_GRAY)
        y += Inches(1.2)

    # ---- Slide 4: Product Overview ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "Product Overview", font_size=32, bold=True, color=DARK_BLUE)
    modules = [
        ("CloudSync Connect", "Multi-provider API gateway with 200+ native integrations"),
        ("CloudSync Shield", "Zero-trust security layer with automated compliance scanning"),
        ("CloudSync Flow", "Workflow automation engine with 50+ pre-built templates"),
        ("CloudSync Insight", "Real-time analytics and cost optimization recommendations"),
    ]
    y = Inches(1.8)
    for mod_name, mod_desc in modules:
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.8), y, Inches(11.5), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
        shape.line.color.rgb = ACCENT_BLUE
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run_title = p.add_run()
        run_title.text = mod_name
        run_title.font.bold = True
        run_title.font.size = Pt(16)
        run_title.font.color.rgb = DARK_BLUE
        run_desc = p.add_run()
        run_desc.text = f"  —  {mod_desc}"
        run_desc.font.size = Pt(14)
        run_desc.font.color.rgb = DARK_GRAY
        y += Inches(1.2)

    # ---- Slide 5: Key Features ----
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "Key Features", font_size=32, bold=True, color=DARK_BLUE)
    kf = [
        "Real-time bi-directional sync across all major cloud providers",
        "Automated compliance reporting for SOC 2, ISO 27001, HIPAA, GDPR",
        "AI-driven cost optimization saving clients an average of 34%",
        "99.99% uptime SLA with global edge network (47 PoPs)",
        "Role-based access control with SSO and MFA integration",
        "Custom API extensions and webhook support for enterprise workflows",
    ]
    y = Inches(1.8)
    for feat in kf:
        add_text_box(slide5, Inches(1.2), y, Inches(10), Inches(0.5),
                     f"✓  {feat}", font_size=15, color=DARK_GRAY)
        y += Inches(0.75)

    # ---- Slide 6: Pricing ----
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "Pricing Plans", font_size=32, bold=True, color=DARK_BLUE)
    plans = [
        ("Starter", "$2,499/mo", "Up to 3 cloud providers, 50 users, email support"),
        ("Professional", "$5,999/mo", "Up to 10 providers, 250 users, priority support, API access"),
        ("Enterprise", "Custom", "Unlimited providers & users, dedicated CSM, custom SLAs, on-prem option"),
    ]
    x = Inches(0.8)
    for plan_name, price, desc in plans:
        shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        x, Inches(1.8), Inches(3.6), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = plan_name
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = DARK_BLUE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = price
        run2.font.size = Pt(28)
        run2.font.bold = True
        run2.font.color.rgb = ACCENT_BLUE
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = desc
        run3.font.size = Pt(12)
        run3.font.color.rgb = LIGHT_GRAY
        x += Inches(4.0)

    # ---- Slide 7: Case Study ----
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    fill7 = slide7.background.fill
    fill7.solid()
    fill7.fore_color.rgb = BG_LIGHT
    add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "Case Study: TechVault Industries", font_size=32, bold=True, color=DARK_BLUE)
    add_text_box(slide7, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
                 "TechVault Industries, a Fortune 500 manufacturing company, deployed CloudSync Pro "
                 "across 12 global offices to unify their hybrid cloud infrastructure.",
                 font_size=16, color=DARK_GRAY)
    metrics = [
        ("62%", "Reduction in cloud management overhead"),
        ("$1.8M", "Annual cost savings in first year"),
        ("99.997%", "Sync reliability across 8 cloud providers"),
        ("3 weeks", "Full deployment and migration timeline"),
    ]
    x = Inches(0.8)
    for val, label in metrics:
        add_text_box(slide7, x, Inches(3.5), Inches(2.8), Inches(0.8),
                     val, font_size=36, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
        add_text_box(slide7, x, Inches(4.5), Inches(2.8), Inches(0.6),
                     label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        x += Inches(3.0)

    # ---- Slide 8: ROI Analysis ----
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
                 "ROI Analysis — Meridian Corp Projection", font_size=32, bold=True, color=DARK_BLUE)
    table_shape = slide8.shapes.add_table(5, 3, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5))
    table = table_shape.table
    headers = ["Category", "Current Annual Cost", "With CloudSync Pro"]
    data_rows = [
        ["Cloud Management Tools", "$1,240,000", "$480,000"],
        ["IT Staff Overhead (Manual Tasks)", "$890,000", "$320,000"],
        ["Compliance & Audit", "$560,000", "$210,000"],
        ["Total", "$2,690,000", "$1,010,000"],
    ]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = WHITE
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        solidFill.append(solidFill.makeelement(qn('a:srgbClr'), {'val': '2B6CB0'}))
        tcPr.append(solidFill)
    for r, row_data in enumerate(data_rows, 1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = val
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(13)
                if r == len(data_rows):
                    run.font.bold = True

    # ---- Slide 9: Q&A ----
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    fill9 = slide9.background.fill
    fill9.solid()
    fill9.fore_color.rgb = DARK_BLUE
    add_text_box(slide9, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
                 "Questions & Discussion", font_size=40, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
                 "We're here to address any questions about how CloudSync Pro\ncan transform your cloud infrastructure.",
                 font_size=18, color=RGBColor(0xBB, 0xD5, 0xED), alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
