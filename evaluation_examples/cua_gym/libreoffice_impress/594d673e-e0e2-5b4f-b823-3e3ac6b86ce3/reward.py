"""
Reward Script: Next Steps call-to-action slide verification
Task ID: impress_sales_078
Domain: libreoffice_impress
Scoring:
  Component 1 (0.15): Slide 10 exists with 'Next Steps' title
  Component 2 (0.35): Three rounded rectangles with correct action item text and #F5F5F5 fill
  Component 3 (0.25): Three blue (#2B6CB0) circles with numbers 1, 2, 3
  Component 4 (0.15): Contact info at bottom with correct content
  Component 5 (0.10): Text alignment correctness (left for items, center for circles/contact)
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_078'

# Expected action items
ACTION_ITEMS = [
    '1. Schedule Technical Deep-Dive (Next Week)',
    '2. Begin Free 30-Day Trial',
    '3. Executive Alignment Meeting',
]

CONTACT_TEXT = 'Your Account Executive: Alex Rivera | alex@cloudsync.io | (555) 987-6543'


def normalize(text):
    """Normalize text for comparison: strip, collapse whitespace."""
    if text is None:
        return ''
    return ' '.join(str(text).strip().split())


def get_fill_hex(shape):
    """Get the solid fill color hex of a shape, or None."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    return None


def get_shape_text(shape):
    """Get concatenated text from a shape's text frame."""
    if not hasattr(shape, 'text_frame'):
        return ''
    return shape.text_frame.text.strip()


def get_para_alignment(para):
    """Get paragraph alignment, treating None as LEFT."""
    align = para.alignment
    if align is None:
        return PP_ALIGN.LEFT
    return align


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    num_slides = len(prs.slides)

    # =====================================================================
    # Component 1: Slide 10 exists with 'Next Steps' title (0.15 points)
    # =====================================================================
    try:
        if num_slides < 10:
            print(f"FAIL: Component 1 — Only {num_slides} slides, need at least 10")
        else:
            slide10 = prs.slides[9]
            # Find a shape with 'Next Steps' text
            title_found = False
            for shape in slide10.shapes:
                text = get_shape_text(shape)
                if 'next steps' in normalize(text).lower():
                    title_found = True
                    break
            if title_found:
                print(f"PASS: Component 1 — Slide 10 exists with 'Next Steps' title (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 1 — Slide 10 exists but no 'Next Steps' title found")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Early exit if no slide 10
    if num_slides < 10:
        final_score = min(total_score, 1.0)
        print(f"\nScore: {total_score}/1.0")
        print(f"REWARD: {final_score}")
        return final_score

    slide10 = prs.slides[9]

    # =====================================================================
    # Component 2: Three rounded rectangles with action item text and
    #              #F5F5F5 fill (0.35 points)
    # =====================================================================
    try:
        # Collect auto shapes that are rounded rectangles with action text
        matched_items = 0
        for item_text in ACTION_ITEMS:
            item_norm = normalize(item_text).lower()
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    shape_text = normalize(get_shape_text(shape)).lower()
                    if item_norm == shape_text or item_norm in shape_text:
                        # Check fill color
                        fill_hex = get_fill_hex(shape)
                        if fill_hex == 'F5F5F5':
                            matched_items += 1
                            print(f"  MATCH: Action item '{item_text}' in shape with #F5F5F5 fill")
                        else:
                            print(f"  PARTIAL: Action item '{item_text}' found but fill={fill_hex}, expected F5F5F5")
                            matched_items += 0.5  # partial credit for text without fill
                        break
            else:
                print(f"  MISS: Action item '{item_text}' not found in any auto shape")

        # Score: proportional to how many matched (3 items = full credit)
        item_score = (matched_items / 3.0) * 0.35
        if matched_items >= 3:
            print(f"PASS: Component 2 — All 3 action items in rounded rectangles with #F5F5F5 fill ({item_score:.2f} pts)")
        elif matched_items > 0:
            print(f"PARTIAL: Component 2 — {matched_items}/3 action items matched ({item_score:.2f} pts)")
        else:
            print(f"FAIL: Component 2 — No action items found in rounded rectangles")
        total_score += item_score
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # =====================================================================
    # Component 3: Three blue (#2B6CB0) circles with numbers 1, 2, 3
    #              (0.25 points)
    # =====================================================================
    try:
        matched_circles = 0
        expected_numbers = ['1', '2', '3']

        for num in expected_numbers:
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    shape_text = get_shape_text(shape).strip()
                    if shape_text == num:
                        fill_hex = get_fill_hex(shape)
                        if fill_hex == '2B6CB0':
                            matched_circles += 1
                            print(f"  MATCH: Circle '{num}' with #2B6CB0 fill")
                        else:
                            print(f"  PARTIAL: Circle '{num}' found but fill={fill_hex}, expected 2B6CB0")
                            matched_circles += 0.5
                        break
            else:
                print(f"  MISS: Number circle '{num}' not found")

        circle_score = (matched_circles / 3.0) * 0.25
        if matched_circles >= 3:
            print(f"PASS: Component 3 — All 3 blue number circles found ({circle_score:.2f} pts)")
        elif matched_circles > 0:
            print(f"PARTIAL: Component 3 — {matched_circles}/3 number circles matched ({circle_score:.2f} pts)")
        else:
            print(f"FAIL: Component 3 — No number circles found")
        total_score += circle_score
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # =====================================================================
    # Component 4: Contact info at bottom (0.15 points)
    # =====================================================================
    try:
        contact_found = False
        contact_keywords = ['alex rivera', 'alex@cloudsync.io', '(555) 987-6543']

        for shape in slide10.shapes:
            text = normalize(get_shape_text(shape)).lower()
            if all(kw in text for kw in contact_keywords):
                contact_found = True
                print(f"  MATCH: Contact info found: '{get_shape_text(shape)[:80]}...'")
                break

        if contact_found:
            print(f"PASS: Component 4 — Contact information present (0.15 pts)")
            total_score += 0.15
        else:
            # Check for partial contact info
            partial = 0
            for shape in slide10.shapes:
                text = normalize(get_shape_text(shape)).lower()
                for kw in contact_keywords:
                    if kw in text:
                        partial += 1
                        break
            if partial > 0:
                partial_score = (partial / 3.0) * 0.15
                print(f"PARTIAL: Component 4 — Partial contact info ({partial}/3 keywords) ({partial_score:.2f} pts)")
                total_score += partial_score
            else:
                print(f"FAIL: Component 4 — No contact information found on slide 10")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # =====================================================================
    # Component 5: Text alignment correctness (0.10 points)
    #   - Action items should be left-aligned
    #   - Number circles should be center-aligned
    # =====================================================================
    try:
        alignment_checks = 0
        total_alignment_checks = 0

        # Check action item alignment (left)
        for item_text in ACTION_ITEMS:
            item_norm = normalize(item_text).lower()
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(shape, 'text_frame'):
                    shape_text = normalize(get_shape_text(shape)).lower()
                    if item_norm == shape_text or item_norm in shape_text:
                        total_alignment_checks += 1
                        for para in shape.text_frame.paragraphs:
                            align = get_para_alignment(para)
                            if align == PP_ALIGN.LEFT:
                                alignment_checks += 1
                            else:
                                print(f"  ALIGN FAIL: '{item_text[:30]}...' is {align}, expected LEFT")
                        break

        # Check circle alignment (center)
        for num in ['1', '2', '3']:
            for shape in slide10.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and hasattr(shape, 'text_frame'):
                    if get_shape_text(shape).strip() == num:
                        total_alignment_checks += 1
                        for para in shape.text_frame.paragraphs:
                            align = get_para_alignment(para)
                            if align == PP_ALIGN.CENTER:
                                alignment_checks += 1
                            else:
                                print(f"  ALIGN FAIL: Circle '{num}' is {align}, expected CENTER")
                        break

        if total_alignment_checks > 0:
            align_score = (alignment_checks / total_alignment_checks) * 0.10
            if alignment_checks == total_alignment_checks:
                print(f"PASS: Component 5 — All alignments correct ({align_score:.2f} pts)")
            else:
                print(f"PARTIAL: Component 5 — {alignment_checks}/{total_alignment_checks} alignments correct ({align_score:.2f} pts)")
            total_score += align_score
        else:
            print(f"FAIL: Component 5 — No shapes found to check alignment")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
