"""
Initial Setup: Training workshop deck — 9 slides with regular-weight text, inconsistent title sizes, no underlines
Task ID: osworld_impress_bold_all_title_size_underline_008
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_008'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_slide_content(prs, layout_idx, title_text, title_size_pt, body_paragraphs):
    """
    Add a slide with:
    - title placeholder (with specific non-bold, no-underline font at given size)
    - content placeholder with body_paragraphs (list of strings), all non-bold
    """
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Set title
    title_shape = slide.shapes.title
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.bold = False
    run.font.underline = False
    run.font.size = Pt(title_size_pt)

    # Set body content
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        body_tf = body_shape.text_frame
        body_tf.clear()
        for i, para_text in enumerate(body_paragraphs):
            if i == 0:
                p = body_tf.paragraphs[0]
            else:
                p = body_tf.add_paragraph()
            p.level = 0
            run = p.add_run()
            run.text = para_text
            run.font.bold = False
            run.font.size = Pt(18)

    return slide


def create_initial():
    prs = Presentation()
    # Use layout 1 (Title and Content) for most slides, 0 for title slide
    # Title slide layout index 0, content slide layout index 1

    # Slide 1 — Title slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    tf1 = title1.text_frame
    tf1.clear()
    p1 = tf1.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Effective Communication in the Workplace"
    r1.font.bold = False
    r1.font.underline = False
    r1.font.size = Pt(32)   # intentionally NOT 30pt

    if len(slide1.placeholders) > 1:
        sub1 = slide1.placeholders[1]
        stf = sub1.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = "Professional Development Workshop Series"
        sr.font.bold = False
        sr.font.size = Pt(18)

    # Slide 2 — Agenda (title size 26pt)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    t2 = slide2.shapes.title
    tf2 = t2.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Workshop Agenda"
    r2.font.bold = False
    r2.font.size = Pt(26)

    body2 = slide2.placeholders[1]
    btf2 = body2.text_frame
    btf2.clear()
    agenda_items = [
        "Introduction and Icebreakers",
        "Core Communication Principles",
        "Active Listening Techniques",
        "Non-Verbal Communication",
        "Conflict Resolution Strategies",
        "Practical Exercises",
        "Q&A and Wrap-Up",
    ]
    for i, item in enumerate(agenda_items):
        if i == 0:
            p = btf2.paragraphs[0]
        else:
            p = btf2.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 3 — Section: Core Principles (title size 38pt)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    t3 = slide3.shapes.title
    tf3 = t3.text_frame
    tf3.clear()
    p3 = tf3.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Core Communication Principles"
    r3.font.bold = False
    r3.font.size = Pt(38)

    body3 = slide3.placeholders[1]
    btf3 = body3.text_frame
    btf3.clear()
    principles = [
        "Clarity: Express ideas simply and directly",
        "Consistency: Align verbal and non-verbal messages",
        "Empathy: Understand the listener's perspective",
        "Feedback: Encourage two-way dialogue",
        "Respect: Value every voice in the conversation",
    ]
    for i, item in enumerate(principles):
        if i == 0:
            p = btf3.paragraphs[0]
        else:
            p = btf3.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 4 — Active Listening (title size 22pt)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    t4 = slide4.shapes.title
    tf4 = t4.text_frame
    tf4.clear()
    p4 = tf4.paragraphs[0]
    r4 = p4.add_run()
    r4.text = "Active Listening Techniques"
    r4.font.bold = False
    r4.font.size = Pt(22)

    body4 = slide4.placeholders[1]
    btf4 = body4.text_frame
    btf4.clear()
    listening = [
        "Maintain appropriate eye contact",
        "Avoid interrupting the speaker",
        "Use open body language",
        "Paraphrase to confirm understanding",
        "Ask clarifying questions thoughtfully",
        "Minimize distractions during conversations",
    ]
    for i, item in enumerate(listening):
        if i == 0:
            p = btf4.paragraphs[0]
        else:
            p = btf4.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 5 — Non-Verbal Communication (title size 34pt)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    t5 = slide5.shapes.title
    tf5 = t5.text_frame
    tf5.clear()
    p5 = tf5.paragraphs[0]
    r5 = p5.add_run()
    r5.text = "Non-Verbal Communication"
    r5.font.bold = False
    r5.font.size = Pt(34)

    body5 = slide5.placeholders[1]
    btf5 = body5.text_frame
    btf5.clear()
    nonverbal = [
        "Facial expressions convey 55% of all communication",
        "Tone of voice accounts for 38% of meaning",
        "Actual words represent only 7% of impact",
        "Posture signals confidence and engagement",
        "Gestures reinforce or contradict spoken words",
        "Proxemics: physical distance affects comfort levels",
    ]
    for i, item in enumerate(nonverbal):
        if i == 0:
            p = btf5.paragraphs[0]
        else:
            p = btf5.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 6 — Conflict Resolution (title size 28pt)
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    t6 = slide6.shapes.title
    tf6 = t6.text_frame
    tf6.clear()
    p6 = tf6.paragraphs[0]
    r6 = p6.add_run()
    r6.text = "Conflict Resolution Strategies"
    r6.font.bold = False
    r6.font.size = Pt(28)

    body6 = slide6.placeholders[1]
    btf6 = body6.text_frame
    btf6.clear()
    conflict = [
        "Identify the root cause, not just the symptom",
        "Focus on interests, not positions",
        "Seek win-win outcomes where possible",
        "Stay calm and regulate emotional responses",
        "Use 'I' statements to avoid blame language",
        "Involve a neutral mediator when needed",
    ]
    for i, item in enumerate(conflict):
        if i == 0:
            p = btf6.paragraphs[0]
        else:
            p = btf6.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 7 — Practical Exercise (title size 24pt)
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    t7 = slide7.shapes.title
    tf7 = t7.text_frame
    tf7.clear()
    p7 = tf7.paragraphs[0]
    r7 = p7.add_run()
    r7.text = "Practical Exercise: Role Play"
    r7.font.bold = False
    r7.font.size = Pt(24)

    body7 = slide7.placeholders[1]
    btf7 = body7.text_frame
    btf7.clear()
    exercise = [
        "Pair up with a colleague for this activity",
        "Scenario A: Delivering difficult feedback",
        "Scenario B: Resolving a misunderstanding",
        "Scenario C: Negotiating a project deadline",
        "Debrief: Discuss what worked and what to improve",
        "Duration: 15 minutes per scenario",
    ]
    for i, item in enumerate(exercise):
        if i == 0:
            p = btf7.paragraphs[0]
        else:
            p = btf7.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 8 — Digital Communication (title size 36pt)
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    t8 = slide8.shapes.title
    tf8 = t8.text_frame
    tf8.clear()
    p8 = tf8.paragraphs[0]
    r8 = p8.add_run()
    r8.text = "Digital Communication Best Practices"
    r8.font.bold = False
    r8.font.size = Pt(36)

    body8 = slide8.placeholders[1]
    btf8 = body8.text_frame
    btf8.clear()
    digital = [
        "Choose the right channel: email vs. chat vs. call",
        "Respond promptly to maintain professional trust",
        "Keep emails concise; use clear subject lines",
        "Avoid tone misinterpretation — use emojis carefully",
        "Check for grammar and clarity before sending",
        "Respect colleagues' off-hours boundaries",
    ]
    for i, item in enumerate(digital):
        if i == 0:
            p = btf8.paragraphs[0]
        else:
            p = btf8.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    # Slide 9 — Closing & Key Takeaways (title size 30pt — one slide already 30pt is fine)
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    t9 = slide9.shapes.title
    tf9 = t9.text_frame
    tf9.clear()
    p9 = tf9.paragraphs[0]
    r9 = p9.add_run()
    r9.text = "Key Takeaways and Next Steps"
    r9.font.bold = False
    r9.font.size = Pt(30)

    body9 = slide9.placeholders[1]
    btf9 = body9.text_frame
    btf9.clear()
    takeaways = [
        "Communication is a skill that improves with practice",
        "Listen more than you speak in most interactions",
        "Be mindful of non-verbal signals you are sending",
        "Approach conflicts as opportunities for growth",
        "Next session: Advanced Presentation Skills — March 20",
        "Feedback forms are available at the exit",
    ]
    for i, item in enumerate(takeaways):
        if i == 0:
            p = btf9.paragraphs[0]
        else:
            p = btf9.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.bold = False
        run.font.size = Pt(16)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup: open file in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
