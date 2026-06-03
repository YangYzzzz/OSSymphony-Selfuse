"""
Reward Script: Bold all text, standardize title sizes to 30pt, underline slide 1 title.
Task ID: osworld_impress_bold_all_title_size_underline_008
Domain: libreoffice_impress
Scoring:
  Component 1: All text runs bold across all 9 slides (0.4 pts)
  Component 2: All title shapes set to 30pt font size (0.3 pts)
  Component 3: Slide 1 title underlined (0.3 pts)
"""

import os
from pptx import Presentation
from pptx.util import Pt

WORKDIR = '/home/user'
TASK_ID = 'osworld_impress_bold_all_title_size_underline_008'

TARGET_TITLE_SIZE_PT = 30.0  # all titles standardized to 30pt
TARGET_TITLE_SIZE_EMU = int(TARGET_TITLE_SIZE_PT * 12700)  # EMU units

def get_all_runs_in_presentation(prs):
    """Return a list of (slide_idx, shape_name, run) for every non-empty run."""
    all_runs = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.text or "").strip():
                            all_runs.append((slide_idx, shape.name, run))
    return all_runs


def get_title_shapes(prs):
    """Return a list of (slide_idx, shape) for title-named shapes on each slide."""
    title_shapes = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name.startswith("Title"):
                title_shapes.append((slide_idx, shape))
    return title_shapes


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    # Precondition: load the file
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: ensure 9 slides exist
    if len(prs.slides) != 9:
        print(f"CRITICAL: Expected 9 slides, found {len(prs.slides)}")
        print("REWARD: 0.0")
        return 0.0

    # -----------------------------------------------------------------------
    # Component 1: All text runs across all slides are bold (0.4 pts)
    # Initial env: all runs have bold=False — this FAILS on initial
    # Golden env: all runs have bold=True  — this PASSES on golden
    # -----------------------------------------------------------------------
    try:
        all_runs = get_all_runs_in_presentation(prs)
        non_bold_runs = []
        for slide_idx, shape_name, run in all_runs:
            bold_val = run.font.bold
            # None means inherited — treat as not explicitly bold
            if bold_val is not True:
                non_bold_runs.append(
                    f"slide {slide_idx+1} shape '{shape_name}' run '{run.text[:30]}'"
                )

        if not non_bold_runs:
            total_runs = len(all_runs)
            print(f"PASS: Component 1 — All {total_runs} text runs are bold (0.4 pts)")
            total_score += 0.4
        else:
            print(
                f"FAIL: Component 1 — {len(non_bold_runs)} run(s) are NOT bold. "
                f"First failures: {non_bold_runs[:3]}"
            )
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # -----------------------------------------------------------------------
    # Component 2: All title shapes have font size exactly 30pt (0.3 pts)
    # Initial env: titles vary from 22pt to 38pt (8 of 9 are not 30pt)
    # Golden env: all titles are exactly 30pt
    # -----------------------------------------------------------------------
    try:
        title_shapes = get_title_shapes(prs)
        wrong_size_titles = []
        for slide_idx, shape in title_shapes:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        sz_emu = run.font.size
                        if sz_emu is None:
                            wrong_size_titles.append(
                                f"slide {slide_idx+1} title: size is None (not explicitly set)"
                            )
                        else:
                            sz_pt = sz_emu / 12700.0
                            if abs(sz_pt - TARGET_TITLE_SIZE_PT) > 0.1:
                                wrong_size_titles.append(
                                    f"slide {slide_idx+1} title: size={sz_pt}pt (expected {TARGET_TITLE_SIZE_PT}pt)"
                                )

        if not wrong_size_titles:
            print(
                f"PASS: Component 2 — All {len(title_shapes)} title shapes are "
                f"{TARGET_TITLE_SIZE_PT}pt (0.3 pts)"
            )
            total_score += 0.3
        else:
            print(
                f"FAIL: Component 2 — {len(wrong_size_titles)} title(s) are NOT {TARGET_TITLE_SIZE_PT}pt. "
                f"Issues: {wrong_size_titles}"
            )
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # -----------------------------------------------------------------------
    # Component 3: Slide 1 title is underlined (0.3 pts)
    # Initial env: slide 1 title underline=False — FAILS on initial
    # Golden env: slide 1 title underline=True  — PASSES on golden
    # -----------------------------------------------------------------------
    try:
        slide1 = prs.slides[0]
        slide1_title_shape = None
        for shape in slide1.shapes:
            if shape.has_text_frame and shape.name.startswith("Title"):
                slide1_title_shape = shape
                break

        if slide1_title_shape is None:
            print("FAIL: Component 3 — No title shape found on slide 1")
        else:
            underlined_runs = []
            total_title_runs = []
            for para in slide1_title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        total_title_runs.append(run)
                        if run.font.underline is True:
                            underlined_runs.append(run)

            if total_title_runs and len(underlined_runs) == len(total_title_runs):
                print(
                    f"PASS: Component 3 — Slide 1 title is underlined "
                    f"({len(underlined_runs)}/{len(total_title_runs)} runs) (0.3 pts)"
                )
                total_score += 0.3
            else:
                found_underlines = [r.font.underline for r in total_title_runs]
                print(
                    f"FAIL: Component 3 — Slide 1 title underline not fully set. "
                    f"Run underline values: {found_underlines}"
                )
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path in the VM env
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
