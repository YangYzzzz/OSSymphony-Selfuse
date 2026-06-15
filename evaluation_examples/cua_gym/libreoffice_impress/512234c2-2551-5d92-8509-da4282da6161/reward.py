"""
Reward Script: Venn diagram with three overlapping circles on slide 6
Task ID: impress_rp_037
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Three ellipse/oval shapes exist on slide 6
  Component 2 (0.25): Correct fill colors (#FF69B4, #4169E1, #32CD32)
  Component 3 (0.25): Correct text labels on circles (Design, Engineering, Business)
  Component 4 (0.25): 'Innovation' text box exists with bold formatting
"""

import os
import zipfile
import xml.etree.ElementTree as ET

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_037'


def persist_app_state(domain):
    """Save any unsaved LibreOffice changes before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print(f"PERSIST: ctrl+s sent for {domain}")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify Venn diagram creation on slide 6 with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Presentation must have at least 6 slides
    if len(prs.slides) < 6:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 6")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[5]  # 0-indexed, slide 6

    # Collect all shapes on slide 6, excluding placeholders (title etc.)
    # We need to find oval/ellipse shapes and text boxes
    non_placeholder_shapes = []
    for shape in slide.shapes:
        if shape.shape_type != 14:  # 14 = PLACEHOLDER
            non_placeholder_shapes.append(shape)

    # Also parse XML to get precise shape geometry info
    ns = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }

    oval_shapes = []
    textbox_shapes = []

    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            with zf.open('ppt/slides/slide6.xml') as f:
                root = ET.fromstring(f.read())

            for sp in root.findall('.//p:cSld/p:spTree/p:sp', ns):
                spPr = sp.find('p:spPr', ns)
                if spPr is None:
                    continue

                prstGeom = spPr.find('a:prstGeom', ns)
                shape_type = prstGeom.get('prst') if prstGeom is not None else None

                # Get text content
                texts = []
                txBody = sp.find('p:txBody', ns)
                if txBody is not None:
                    for p_elem in txBody.findall('a:p', ns):
                        for r_elem in p_elem.findall('a:r', ns):
                            t = r_elem.find('a:t', ns)
                            if t is not None and t.text:
                                texts.append(t.text.strip())

                # Get fill color
                fill_color = None
                solidFill = spPr.find('a:solidFill', ns)
                if solidFill is not None:
                    srgb = solidFill.find('a:srgbClr', ns)
                    if srgb is not None:
                        fill_color = srgb.get('val')

                # Get bold info by checking run properties
                bold_runs = []
                if txBody is not None:
                    for p_elem in txBody.findall('a:p', ns):
                        for r_elem in p_elem.findall('a:r', ns):
                            rPr = r_elem.find('a:rPr', ns)
                            if rPr is not None and rPr.get('b') == '1':
                                bold_runs.append(rPr)
                is_bold = len(bold_runs) > 0

                # Check for placeholder (skip)
                nvSpPr = sp.find('p:nvSpPr', ns)
                nvPr = nvSpPr.find('p:nvPr', ns) if nvSpPr is not None else None
                ph = nvPr.find('p:ph', ns) if nvPr is not None else None
                if ph is not None:
                    continue

                shape_info = {
                    'type': shape_type,
                    'texts': texts,
                    'fill_color': fill_color,
                    'is_bold': is_bold,
                }

                if shape_type == 'ellipse':
                    oval_shapes.append(shape_info)
                elif shape_type == 'rect' or shape_type is None:
                    # Text boxes often have rect or no preset geometry
                    if texts:
                        textbox_shapes.append(shape_info)
    except Exception as e:
        print(f"ERROR: XML parsing failed: {e}")
        print("REWARD: 0.0")
        return 0.0

    print(f"INFO: Found {len(oval_shapes)} ellipse shapes and {len(textbox_shapes)} text-bearing non-ellipse shapes on slide 6")

    # Component 1: Three ellipse/oval shapes exist on slide 6 (0.25 points)
    # This tests the core structural change - adding 3 circle shapes
    try:
        if len(oval_shapes) >= 3:
            print(f"PASS: Component 1 — Found {len(oval_shapes)} ellipse shapes on slide 6 (0.25 pts)")
            total_score += 0.25
        else:
            print(f"FAIL: Component 1 — Expected 3 ellipse shapes, found {len(oval_shapes)}")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Correct fill colors (#FF69B4, #4169E1, #32CD32) (0.25 points)
    # Each correct color earns ~0.083 points
    try:
        expected_colors = {'FF69B4', '4169E1', '32CD32'}
        found_colors = set()
        for oval in oval_shapes:
            if oval['fill_color'] and oval['fill_color'].upper() in {c.upper() for c in expected_colors}:
                found_colors.add(oval['fill_color'].upper())

        color_score = len(found_colors) / 3.0 * 0.25
        if len(found_colors) == 3:
            print(f"PASS: Component 2 — All 3 colors found: {found_colors} (0.25 pts)")
            total_score += 0.25
        elif len(found_colors) > 0:
            print(f"PARTIAL: Component 2 — Found {len(found_colors)}/3 colors: {found_colors} ({color_score:.3f} pts)")
            total_score += color_score
        else:
            actual_colors = [oval['fill_color'] for oval in oval_shapes]
            print(f"FAIL: Component 2 — No matching colors found. Actual: {actual_colors}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Correct text labels on circles (Design, Engineering, Business) (0.25 points)
    # Each correct label earns ~0.083 points
    try:
        expected_labels = {'design', 'engineering', 'business'}
        found_labels = set()
        # Map: color -> expected label
        color_label_map = {
            'FF69B4': 'design',
            '4169E1': 'engineering',
            '32CD32': 'business',
        }

        for oval in oval_shapes:
            oval_text = ' '.join(oval['texts']).strip().lower()
            if oval_text in expected_labels:
                found_labels.add(oval_text)
                # Bonus: check color-label association
                if oval['fill_color']:
                    expected_label = color_label_map.get(oval['fill_color'].upper(), '')
                    if expected_label == oval_text:
                        pass  # correct association

        label_score = len(found_labels) / 3.0 * 0.25
        if len(found_labels) == 3:
            print(f"PASS: Component 3 — All 3 labels found: {found_labels} (0.25 pts)")
            total_score += 0.25
        elif len(found_labels) > 0:
            print(f"PARTIAL: Component 3 — Found {len(found_labels)}/3 labels: {found_labels} ({label_score:.3f} pts)")
            total_score += label_score
        else:
            actual_texts = [oval['texts'] for oval in oval_shapes]
            print(f"FAIL: Component 3 — No matching labels. Actual texts: {actual_texts}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 'Innovation' text box in center overlap area with bold formatting (0.25 points)
    try:
        # Check both textbox shapes and oval shapes for 'Innovation' text
        all_text_shapes = textbox_shapes + oval_shapes
        innovation_matches = [s for s in all_text_shapes if 'innovation' in ' '.join(s['texts']).strip().lower()]
        innovation_bold_matches = [s for s in innovation_matches if s['is_bold']]

        if len(innovation_bold_matches) > 0:
            print(f"PASS: Component 4 — 'Innovation' text found with bold formatting (0.25 pts)")
            total_score += 0.25
        elif len(innovation_matches) > 0:
            print(f"PARTIAL: Component 4 — 'Innovation' text found but NOT bold (0.125 pts)")
            total_score += 0.125
        else:
            print(f"FAIL: Component 4 — 'Innovation' text not found on slide 6")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
