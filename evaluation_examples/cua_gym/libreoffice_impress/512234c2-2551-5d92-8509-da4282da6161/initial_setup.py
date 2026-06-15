"""
Initial Setup: Create a 10-slide Cross-Functional presentation with slide 6 blank (title only)
Task ID: impress_rp_037
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_rp_037'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_content(slide, title_text, bullets):
    """Helper to populate a title+content slide."""
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        if i == 0:
            tf.paragraphs[0].text = text
        else:
            p = tf.add_paragraph()
            p.text = text


def add_blank_with_title(prs, title_text):
    """Add a slide with layout 6 (Title Only) and set the title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.title.text = title_text
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Cross-Functional Collaboration"
    slide1.placeholders[1].text = "Strategies for Integrated Product Development\nQ2 2025 Planning Session"

    # --- Slide 2: Agenda ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide2, "Agenda", [
        "Overview of Cross-Functional Teams",
        "Key Challenges in 2024",
        "Proposed Framework for Collaboration",
        "Innovation at the Intersection",
        "Case Studies & Metrics",
        "Next Steps and Action Items",
    ])

    # --- Slide 3: Team Overview ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide3, "Our Cross-Functional Teams", [
        "Design: 14 members across UX, Visual, and Research",
        "Engineering: 32 members spanning Frontend, Backend, and DevOps",
        "Business: 18 members in Product, Marketing, and Sales",
        "Data Science: 8 analysts embedded across all groups",
    ])

    # --- Slide 4: Challenges ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide4, "Key Challenges in Cross-Functional Work", [
        "Siloed decision-making reduces agility",
        "Misaligned priorities between departments",
        "Communication overhead in large teams",
        "Difficulty measuring shared outcomes",
        "Tool fragmentation across disciplines",
    ])

    # --- Slide 5: Framework ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide5, "Proposed Collaboration Framework", [
        "Shared OKRs across Design, Engineering, and Business",
        "Weekly cross-team standups (Tuesdays, 10 AM)",
        "Unified project tracker in Jira with shared boards",
        "Monthly innovation sprints for experimental projects",
        "Quarterly retrospectives with all stakeholders",
    ])

    # --- Slide 6: Where Innovation Happens (BLANK - title only) ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
    slide6.shapes.title.text = "Where Innovation Happens"
    # This slide is intentionally left blank aside from the title.
    # The task is to add a Venn diagram here.

    # --- Slide 7: Case Study ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide7, "Case Study: Project Aurora", [
        "Goal: Redesign onboarding flow to improve conversion by 25%",
        "Team: 3 designers, 5 engineers, 2 product managers",
        "Timeline: 8 weeks from ideation to launch",
        "Result: 31% improvement in trial-to-paid conversion",
        "Key insight: Co-located teams shipped 2x faster",
    ])

    # --- Slide 8: Metrics ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide8, "Collaboration Metrics Dashboard", [
        "Cross-team PRs merged: 147 (up 42% from Q1)",
        "Joint design reviews completed: 23 sessions",
        "Average feature cycle time: 3.2 weeks (down from 5.1)",
        "Shared Slack channels active: 12 channels, 340 members",
        "NPS for internal collaboration: 72 (target: 65)",
    ])

    # --- Slide 9: Innovation Pipeline ---
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide9, "Innovation Pipeline for Q3", [
        "AI-powered customer segmentation tool (Design + Data Science)",
        "Real-time collaboration canvas for remote teams (Engineering + Design)",
        "Predictive churn model integrated into CRM (Business + Data Science)",
        "Automated accessibility audit system (Engineering + Design)",
    ])

    # --- Slide 10: Next Steps ---
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_content(slide10, "Next Steps & Action Items", [
        "Finalize shared OKR document by April 15, 2025",
        "Launch first cross-team innovation sprint: May 5-16",
        "Set up unified Jira boards for all active projects",
        "Schedule quarterly retrospective for June 30",
        "Assign cross-functional ambassadors per department",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
