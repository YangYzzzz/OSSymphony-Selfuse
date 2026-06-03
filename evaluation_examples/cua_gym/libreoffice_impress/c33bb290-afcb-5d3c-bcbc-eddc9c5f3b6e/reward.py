"""
Reward Script: Verify hexagon differentiator shapes on slide 4
Task ID: impress_sales_092
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Five hexagon AUTO_SHAPE shapes on slide 4
  Component 2 (0.25): Correct text labels on hexagons
  Component 3 (0.25): Correct fill colors on hexagons
  Component 4 (0.25): White text color and drop shadows on hexagons
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_092'

# Expected hexagon labels and their fill colors
EXPECTED_HEXAGONS = {
    'AI-Powered': '2B6CB0',
    'Real-Time': '4CAF50',
    'Scalable': 'FF6B35',
    'Secure': '9C27B0',
    'Integrated': 'F44336',
}


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation has at least 4 slides
    if len(prs.slides) < 4:
        print(f"CRITICAL: Only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[3]  # Slide 4 (0-indexed)

    # Collect hexagon shapes (AUTO_SHAPE type with HEXAGON auto_shape_type)
    hexagons = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.HEXAGON:
                    hexagons.append(shape)
                else:
                    # Accept any auto shape that has text matching expected labels
                    text = shape.text.strip()
                    if text in EXPECTED_HEXAGONS:
                        hexagons.append(shape)
            except Exception:
                # Fallback: accept auto shapes with matching text
                text = shape.text.strip()
                if text in EXPECTED_HEXAGONS:
                    hexagons.append(shape)

    # Component 1: Five hexagon shapes exist on slide 4 (0.25 points)
    try:
        num_hexagons = len(hexagons)
        if num_hexagons >= 5:
            print(f"PASS: Component 1 -- Found {num_hexagons} hexagon shapes on slide 4 (0.25 pts)")
            total_score += 0.25
        elif num_hexagons >= 3:
            partial = 0.25 * (num_hexagons / 5.0)
            print(f"PARTIAL: Component 1 -- Found {num_hexagons}/5 hexagon shapes ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Found {num_hexagons}/5 hexagon shapes")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Build a mapping of text -> shape for found hexagons
    hex_by_text = {}
    for h in hexagons:
        text = h.text.strip()
        hex_by_text[text] = h

    # Component 2: Correct text labels (0.25 points)
    try:
        matched_labels = 0
        for label in EXPECTED_HEXAGONS:
            if label in hex_by_text:
                matched_labels += 1
                print(f"  Label '{label}': FOUND")
            else:
                print(f"  Label '{label}': MISSING")

        if matched_labels == 5:
            print(f"PASS: Component 2 -- All 5 labels present (0.25 pts)")
            total_score += 0.25
        elif matched_labels > 0:
            partial = 0.25 * (matched_labels / 5.0)
            print(f"PARTIAL: Component 2 -- {matched_labels}/5 labels present ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 2 -- No expected labels found")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Correct fill colors (0.25 points)
    try:
        correct_colors = 0
        for label, expected_color in EXPECTED_HEXAGONS.items():
            if label not in hex_by_text:
                print(f"  Color for '{label}': SKIP (shape not found)")
                continue
            shape = hex_by_text[label]
            fill = shape.fill
            if fill.type is not None:
                try:
                    actual_color = str(fill.fore_color.rgb).upper()
                    expected_upper = expected_color.upper()
                    if actual_color == expected_upper:
                        correct_colors += 1
                        print(f"  Color for '{label}': MATCH ({actual_color})")
                    else:
                        print(f"  Color for '{label}': MISMATCH (expected {expected_upper}, got {actual_color})")
                except Exception as e:
                    print(f"  Color for '{label}': ERROR reading rgb -- {e}")
            else:
                print(f"  Color for '{label}': NO FILL")

        if correct_colors == 5:
            print(f"PASS: Component 3 -- All 5 fill colors correct (0.25 pts)")
            total_score += 0.25
        elif correct_colors > 0:
            partial = 0.25 * (correct_colors / 5.0)
            print(f"PARTIAL: Component 3 -- {correct_colors}/5 fill colors correct ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 3 -- No correct fill colors")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: White text color AND drop shadows (0.25 points)
    # Sub-component 4a: White text (0.125 points)
    # Sub-component 4b: Drop shadows (0.125 points)
    try:
        white_text_count = 0
        shadow_count = 0

        for label in EXPECTED_HEXAGONS:
            if label not in hex_by_text:
                continue
            shape = hex_by_text[label]

            # Check white text
            has_white = False
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.type is not None:
                            rgb_str = str(run.font.color.rgb).upper()
                            if rgb_str == 'FFFFFF':
                                has_white = True
                    except Exception:
                        pass
            if has_white:
                white_text_count += 1
                print(f"  White text for '{label}': YES")
            else:
                print(f"  White text for '{label}': NO")

            # Check drop shadow via XML
            sp = shape._element
            effectLst = sp.find('.//' + qn('a:effectLst'))
            has_shadow = False
            if effectLst is not None:
                outerShdw = effectLst.find(qn('a:outerShdw'))
                if outerShdw is not None:
                    has_shadow = True
            if has_shadow:
                shadow_count += 1
                print(f"  Drop shadow for '{label}': YES")
            else:
                print(f"  Drop shadow for '{label}': NO")

        # Score white text sub-component (0.125)
        white_score = 0.0
        if white_text_count == 5:
            white_score = 0.125
            print(f"PASS: Component 4a -- All 5 hexagons have white text (0.125 pts)")
        elif white_text_count > 0:
            white_score = 0.125 * (white_text_count / 5.0)
            print(f"PARTIAL: Component 4a -- {white_text_count}/5 have white text ({white_score:.3f} pts)")
        else:
            print(f"FAIL: Component 4a -- No hexagons have white text")

        # Score drop shadow sub-component (0.125)
        shadow_score = 0.0
        if shadow_count == 5:
            shadow_score = 0.125
            print(f"PASS: Component 4b -- All 5 hexagons have drop shadows (0.125 pts)")
        elif shadow_count > 0:
            shadow_score = 0.125 * (shadow_count / 5.0)
            print(f"PARTIAL: Component 4b -- {shadow_count}/5 have drop shadows ({shadow_score:.3f} pts)")
        else:
            print(f"FAIL: Component 4b -- No hexagons have drop shadows")

        total_score += white_score + shadow_score

    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
