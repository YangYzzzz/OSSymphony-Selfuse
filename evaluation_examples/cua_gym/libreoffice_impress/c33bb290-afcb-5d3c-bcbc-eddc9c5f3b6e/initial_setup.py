"""
Initial Setup: Create a sales pitch presentation with 8 slides.
Slide 4 has the title 'What Sets Us Apart' with empty content area.
Task ID: impress_sales_092
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_092'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_body_slide(prs, title_text, bullets, layout_idx=1):
    """Add a slide with title and bulleted body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(18)
    return slide


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Differentiator Pitch"
    slide1.placeholders[1].text = "NextGen Solutions Inc.\nQ2 2025 Sales Strategy"

    # --- Slide 2: Company Overview ---
    add_title_body_slide(prs, "Company Overview", [
        "Founded in 2018, headquartered in San Francisco",
        "Over 350 enterprise clients across 40 countries",
        "Annual recurring revenue of $87.5M in FY2024",
        "Named a Leader in Gartner Magic Quadrant 2024",
        "Team of 520+ engineers and data scientists",
    ])

    # --- Slide 3: Market Opportunity ---
    add_title_body_slide(prs, "Market Opportunity", [
        "Total addressable market: $42B by 2027 (IDC)",
        "Enterprise automation spending up 28% YoY",
        "Key verticals: Financial Services, Healthcare, Manufacturing",
        "Competitive landscape fragmented with no clear leader",
        "Strategic window for platform consolidation play",
    ])

    # --- Slide 4: What Sets Us Apart (EMPTY content area) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    # Add only the title as a text box
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "What Sets Us Apart"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # --- Slide 5: Product Portfolio ---
    add_title_body_slide(prs, "Product Portfolio", [
        "CorePlatform: Unified data integration engine",
        "InsightHub: Real-time analytics dashboard",
        "AutoFlow: Process automation toolkit",
        "SecureVault: Enterprise-grade data protection",
        "ConnectAPI: Third-party integration framework",
    ])

    # --- Slide 6: Client Success Stories ---
    add_title_body_slide(prs, "Client Success Stories", [
        "Meridian Bank: 45% reduction in processing time",
        "HealthFirst Network: $12M annual cost savings",
        "Atlas Manufacturing: 3x improvement in throughput",
        "Pinnacle Retail: 99.97% system uptime achieved",
    ])

    # --- Slide 7: Growth Strategy ---
    add_title_body_slide(prs, "Growth Strategy", [
        "Expand into APAC and EMEA markets by Q4 2025",
        "Launch vertical-specific solutions for top 5 industries",
        "Strengthen channel partner ecosystem to 200+ partners",
        "Invest $15M in R&D for next-gen AI capabilities",
        "Target 40% revenue growth through land-and-expand model",
    ])

    # --- Slide 8: Next Steps ---
    add_title_body_slide(prs, "Next Steps & Timeline", [
        "Phase 1 (Q2): Pilot program with 10 strategic accounts",
        "Phase 2 (Q3): Full product launch with partner enablement",
        "Phase 3 (Q4): International expansion and enterprise rollout",
        "Key milestone: 500 enterprise clients by end of FY2025",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
