"""
Reward Script: Create attendance table on slide 2 with specific formatting
Task ID: impress_teach_083
Domain: libreoffice_impress
Scoring:
  Component 1: Table exists on slide 2 with correct dimensions (6x12) — 0.20
  Component 2: Header row text matches spec — 0.15
  Component 3: No. column filled with 1-11 — 0.15
  Component 4: Header row fill color #37474F — 0.20
  Component 5: Header text is white (#FFFFFF) and bold — 0.15
  Component 6: All cells have #BDBDBD ~0.5pt borders — 0.15
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_083'

EXPECTED_HEADERS = ['No.', 'Student Name', 'Week 1', 'Week 2', 'Week 3', 'Week 4']


def find_table_on_slide(slide):
    """Find the first TABLE shape on a slide, return table object or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return shape.table
    return None


def get_cell_fill_color(cell):
    """Extract solid fill color hex from a table cell, or None."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is not None:
        solidFill = tcPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgbClr = solidFill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                return srgbClr.get('val').upper()
    return None


def get_cell_border_info(cell, border_name):
    """Get border width (EMU) and color for a cell border (lnL, lnR, lnT, lnB).
    Returns (width_emu, color_hex) or (None, None)."""
    tc = cell._tc
    tcPr = tc.find(qn('a:tcPr'))
    if tcPr is None:
        return None, None
    ln = tcPr.find(qn('a:' + border_name))
    if ln is None:
        return None, None
    w = ln.get('w')
    clr = None
    fill = ln.find(qn('a:solidFill'))
    if fill is not None:
        srgb = fill.find(qn('a:srgbClr'))
        if srgb is not None:
            clr = srgb.get('val').upper()
    return (int(w) if w else None), clr


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 2 slides
    if len(prs.slides) < 2:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 2")
        print("REWARD: 0.0")
        return 0.0

    slide2 = prs.slides[1]
    table = find_table_on_slide(slide2)

    # Component 1: Table exists on slide 2 with correct dimensions 6 cols x 12 rows (0.20 points)
    try:
        if table is None:
            print("FAIL: Component 1 — No table found on slide 2")
            print("REWARD: 0.0")
            return 0.0  # No table means nothing else to check
        rows = len(table.rows)
        cols = len(table.columns)
        if rows == 12 and cols == 6:
            print(f"PASS: Component 1 — Table is {rows}x{cols} (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 1 — Table is {rows}x{cols}, expected 12x6")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Header row text matches expected headers (0.15 points)
    try:
        header_texts = [table.cell(0, c).text.strip() for c in range(min(6, len(table.columns)))]
        if header_texts == EXPECTED_HEADERS:
            print(f"PASS: Component 2 — Headers match: {header_texts} (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 2 — Headers: {header_texts}, expected: {EXPECTED_HEADERS}")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: No. column filled with 1-11 in rows 1-11 (0.15 points)
    try:
        no_values = []
        correct_count = 0
        for r in range(1, min(12, len(table.rows))):
            val = table.cell(r, 0).text.strip()
            no_values.append(val)
            if val == str(r):
                correct_count += 1

        if correct_count == 11:
            print(f"PASS: Component 3 — No. column has 1-11 correctly (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 3 — {correct_count}/11 correct. Values: {no_values}")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: Header row fill color is #37474F (0.20 points)
    try:
        correct_fill = 0
        for c in range(min(6, len(table.columns))):
            fill_color = get_cell_fill_color(table.cell(0, c))
            if fill_color == '37474F':
                correct_fill += 1
            else:
                print(f"  Header cell [0,{c}] fill: {fill_color}, expected 37474F")

        if correct_fill == 6:
            print(f"PASS: Component 4 — All 6 header cells have #37474F fill (0.20 pts)")
            total_score += 0.20
        else:
            print(f"FAIL: Component 4 — {correct_fill}/6 header cells have correct fill")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    # Component 5: Header text is white (#FFFFFF) and bold (0.15 points)
    try:
        white_bold_count = 0
        for c in range(min(6, len(table.columns))):
            cell = table.cell(0, c)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    is_bold = run.font.bold is True
                    try:
                        color_rgb = str(run.font.color.rgb).upper()
                    except Exception:
                        color_rgb = None
                    if is_bold and color_rgb == 'FFFFFF':
                        white_bold_count += 1
                    else:
                        print(f"  Header [{0},{c}] run: bold={run.font.bold}, color={color_rgb}")

        if white_bold_count >= 6:
            print(f"PASS: Component 5 — All header runs are white+bold (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 5 — {white_bold_count}/6 header runs are white+bold")
    except Exception as e:
        print(f"ERROR: Component 5 — {e}")

    # Component 6: All cells have #BDBDBD borders at ~0.5pt (6350 EMU) (0.15 points)
    try:
        total_borders = 0
        correct_borders = 0
        # Check a representative sample of cells across the table
        sample_cells = [(r, c) for r in [0, 1, 5, 11] for c in [0, 2, 5]]
        for r, c in sample_cells:
            if r < len(table.rows) and c < len(table.columns):
                for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                    total_borders += 1
                    w, clr = get_cell_border_info(table.cell(r, c), border_name)
                    # 0.5pt = 6350 EMU. Allow some tolerance (5000-8000)
                    if w is not None and 5000 <= w <= 8000 and clr == 'BDBDBD':
                        correct_borders += 1

        if total_borders > 0 and correct_borders == total_borders:
            print(f"PASS: Component 6 — {correct_borders}/{total_borders} sampled borders correct (0.15 pts)")
            total_score += 0.15
        else:
            print(f"FAIL: Component 6 — {correct_borders}/{total_borders} sampled borders correct")
    except Exception as e:
        print(f"ERROR: Component 6 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
persist_app_state()

if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
