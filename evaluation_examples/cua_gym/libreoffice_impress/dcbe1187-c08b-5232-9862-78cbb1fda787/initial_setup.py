"""
Initial Setup: Create a 4-slide Classroom Management presentation with empty Slide 2
Task ID: impress_teach_083
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_083'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Classroom Management"
    slide1.placeholders[1].text = "Spring 2025 — Prof. Elena Rodriguez"

    # --- Slide 2: Attendance Tracker (EMPTY - just title) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Attendance Tracker"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    # No table here — agent must create it

    # --- Slide 3: Grading Policy ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox3_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf3t = txBox3_title.text_frame
    p3t = tf3t.paragraphs[0]
    p3t.text = "Grading Policy"
    p3t.alignment = PP_ALIGN.LEFT
    r3t = p3t.runs[0]
    r3t.font.size = Pt(28)
    r3t.font.bold = True
    r3t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    txBox3_body = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(10), Inches(5))
    tf3b = txBox3_body.text_frame
    tf3b.word_wrap = True
    items = [
        ("Participation", "15% — Active engagement in class discussions and group activities"),
        ("Homework", "25% — Weekly problem sets submitted via the online portal"),
        ("Midterm Exam", "25% — Covers Chapters 1-6, scheduled for Week 8"),
        ("Final Project", "35% — Team-based research presentation due in Week 15"),
    ]
    for i, (title, desc) in enumerate(items):
        if i == 0:
            p = tf3b.paragraphs[0]
        else:
            p = tf3b.add_paragraph()
        p.text = f"{title}: {desc}"
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(14)

    # --- Slide 4: Class Schedule ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    txBox4_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf4t = txBox4_title.text_frame
    p4t = tf4t.paragraphs[0]
    p4t.text = "Class Schedule"
    p4t.alignment = PP_ALIGN.LEFT
    r4t = p4t.runs[0]
    r4t.font.size = Pt(28)
    r4t.font.bold = True
    r4t.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    txBox4_body = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(10), Inches(5))
    tf4b = txBox4_body.text_frame
    tf4b.word_wrap = True
    schedule_items = [
        "Monday 9:00 AM — 10:30 AM: Lecture (Room 204, Harris Hall)",
        "Wednesday 9:00 AM — 10:30 AM: Lecture (Room 204, Harris Hall)",
        "Friday 2:00 PM — 3:30 PM: Lab Section (Room B12, Science Building)",
        "Office Hours: Tuesday & Thursday 1:00 PM — 3:00 PM (Office 318)",
    ]
    for i, item in enumerate(schedule_items):
        if i == 0:
            p = tf4b.paragraphs[0]
        else:
            p = tf4b.add_paragraph()
        p.text = item
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(14)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
