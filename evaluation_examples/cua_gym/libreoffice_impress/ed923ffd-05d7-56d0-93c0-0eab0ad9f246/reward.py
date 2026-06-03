"""
Reward Script: Design a complete lecture title slide for 'Advanced Organic Chemistry'
Task ID: impress_teach_054
Domain: libreoffice_impress
Scoring:
  Component 1 (0.25): Gradient background (dark blue #0D47A1 to black #000000)
  Component 2 (0.30): Title text 'Advanced Organic Chemistry' in 48pt white bold
  Component 3 (0.25): Subtitle '2025 Spring Semester | Dr. James Wilson' in 20pt #BDBDBD
  Component 4 (0.20): Gold (#FFD700) horizontal line, 3pt width
"""

import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_054'


def color_close(actual_rgb, expected_hex, tolerance=5):
    """Check if an RGBColor is close to expected hex string (e.g. '0D47A1')."""
    if actual_rgb is None:
        return False
    actual_str = str(actual_rgb).upper()
    expected_hex = expected_hex.upper()
    try:
        ar, ag, ab = int(actual_str[0:2], 16), int(actual_str[2:4], 16), int(actual_str[4:6], 16)
        er, eg, eb = int(expected_hex[0:2], 16), int(expected_hex[2:4], 16), int(expected_hex[4:6], 16)
        return abs(ar - er) <= tolerance and abs(ag - eg) <= tolerance and abs(ab - eb) <= tolerance
    except (ValueError, IndexError):
        return False


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    if len(prs.slides) == 0:
        print("CRITICAL: No slides in presentation")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[0]

    # Component 1: Gradient background from #0D47A1 to #000000 (0.25 points)
    try:
        fill = slide.background.fill
        if fill.type is not None and fill.type == 3:  # GRADIENT
            stops = fill.gradient_stops
            if len(stops) >= 2:
                # Find stops: top color (#0D47A1) at position 0.0, bottom (#000000) at position 1.0
                stop0_color = stops[0].color.rgb
                stop1_color = stops[-1].color.rgb
                has_blue = color_close(stop0_color, '0D47A1')
                has_black = color_close(stop1_color, '000000')
                if has_blue and has_black:
                    print(f"PASS: Component 1 -- Gradient background correct: {stop0_color} -> {stop1_color} (0.25 pts)")
                    total_score += 0.25
                else:
                    print(f"FAIL: Component 1 -- Gradient colors wrong: {stop0_color} -> {stop1_color}, expected 0D47A1 -> 000000")
            else:
                print(f"FAIL: Component 1 -- Gradient has {len(stops)} stops, expected >= 2")
        else:
            print(f"FAIL: Component 1 -- Background fill type is {fill.type}, expected GRADIENT (3)")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Title 'Advanced Organic Chemistry' in 48pt white bold (0.30 points)
    try:
        title_found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text.strip()
            if 'Advanced Organic Chemistry' in full_text:
                title_found = True
                comp2_score = 0.0
                # Check text content
                comp2_score += 0.10
                # Check font properties on runs
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if 'Advanced Organic Chemistry' in run.text:
                            # Check bold
                            if run.font.bold is True:
                                comp2_score += 0.05
                            else:
                                print(f"FAIL: Component 2 sub -- Title bold={run.font.bold}, expected True")
                            # Check size ~48pt (609600 EMU)
                            if run.font.size is not None and abs(run.font.size - Pt(48)) <= Pt(2):
                                comp2_score += 0.05
                            else:
                                print(f"FAIL: Component 2 sub -- Title size={run.font.size}, expected ~{Pt(48)}")
                            # Check white color
                            try:
                                if run.font.color.type is not None and color_close(run.font.color.rgb, 'FFFFFF'):
                                    comp2_score += 0.10
                                else:
                                    actual_c = run.font.color.rgb if run.font.color.type is not None else None
                                    print(f"FAIL: Component 2 sub -- Title color={actual_c}, expected FFFFFF")
                            except Exception:
                                print(f"FAIL: Component 2 sub -- Title color not accessible")
                            break
                    if comp2_score > 0.10:
                        break
                if comp2_score > 0:
                    print(f"PASS: Component 2 -- Title text and properties verified ({comp2_score} pts)")
                    total_score += comp2_score
                break
        if not title_found:
            print("FAIL: Component 2 -- Title 'Advanced Organic Chemistry' not found in any shape")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Subtitle '2025 Spring Semester | Dr. James Wilson' in 20pt #BDBDBD (0.25 points)
    try:
        subtitle_found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = shape.text_frame.text.strip()
            if '2025 Spring Semester' in full_text and 'Dr. James Wilson' in full_text:
                subtitle_found = True
                comp3_score = 0.0
                comp3_score += 0.10  # text content present
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if '2025 Spring Semester' in run.text or 'Dr. James Wilson' in run.text:
                            # Check size ~20pt (254000 EMU)
                            if run.font.size is not None and abs(run.font.size - Pt(20)) <= Pt(2):
                                comp3_score += 0.05
                            else:
                                print(f"FAIL: Component 3 sub -- Subtitle size={run.font.size}, expected ~{Pt(20)}")
                            # Check color #BDBDBD
                            try:
                                if run.font.color.type is not None and color_close(run.font.color.rgb, 'BDBDBD'):
                                    comp3_score += 0.10
                                else:
                                    actual_c = run.font.color.rgb if run.font.color.type is not None else None
                                    print(f"FAIL: Component 3 sub -- Subtitle color={actual_c}, expected BDBDBD")
                            except Exception:
                                print(f"FAIL: Component 3 sub -- Subtitle color not accessible")
                            break
                    if comp3_score > 0.10:
                        break
                if comp3_score > 0:
                    print(f"PASS: Component 3 -- Subtitle text and properties verified ({comp3_score} pts)")
                    total_score += comp3_score
                break
        if not subtitle_found:
            print("FAIL: Component 3 -- Subtitle text not found in any shape")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Gold (#FFD700) horizontal line, 3pt width (0.20 points)
    try:
        line_found = False
        for shape in slide.shapes:
            # Look for LINE shape type (9) or connector shapes
            if shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.shape_type == 9:
                line_found = True
                comp4_score = 0.0
                # Check it's roughly horizontal (height == 0 or very small relative to width)
                if shape.width > 0 and (shape.height == 0 or abs(shape.height) < shape.width * 0.05):
                    comp4_score += 0.05
                else:
                    print(f"FAIL: Component 4 sub -- Line not horizontal: w={shape.width}, h={shape.height}")
                # Check line color is gold (#FFD700)
                try:
                    ln = shape.line
                    if ln.fill.type is not None and ln.fill.type == 1:  # SOLID
                        if color_close(ln.color.rgb, 'FFD700'):
                            comp4_score += 0.10
                        else:
                            print(f"FAIL: Component 4 sub -- Line color={ln.color.rgb}, expected FFD700")
                    else:
                        print(f"FAIL: Component 4 sub -- Line fill type={ln.fill.type}, expected SOLID (1)")
                except Exception as e2:
                    print(f"FAIL: Component 4 sub -- Line color check failed: {e2}")
                # Check line width ~3pt (38100 EMU)
                if ln.width is not None and abs(ln.width - Pt(3)) <= Pt(1):
                    comp4_score += 0.05
                else:
                    print(f"FAIL: Component 4 sub -- Line width={ln.width}, expected ~{Pt(3)}")
                if comp4_score > 0:
                    print(f"PASS: Component 4 -- Gold horizontal line verified ({comp4_score} pts)")
                    total_score += comp4_score
                break
        if not line_found:
            print("FAIL: Component 4 -- No LINE shape found on slide 1")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = round(min(total_score, 1.0), 2)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
