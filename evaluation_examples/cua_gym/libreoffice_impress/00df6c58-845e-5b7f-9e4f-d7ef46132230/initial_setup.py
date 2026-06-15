"""
Initial Setup: Create project_slides.pptx with 8 slides for slide management task
Task ID: impress_gf5_002
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_002'
FILENAME = 'project_slides.pptx'
OUTPUT = f'{WORKDIR}/{FILENAME}'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_and_body(slide, title_text, body_lines):
    """Helper to set title and body content on a slide."""
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    # Find body placeholder (index 1 typically)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(body_lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
            break


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Project Horizon - Q3 Initiative"
    slide1.placeholders[1].text = "Strategic Roadmap & Execution Plan\nPrepared by Sarah Chen, VP of Operations"

    # --- Slide 2: Phase 1 Plan ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide2, "Phase 1 Plan", [
        "Market Research & Competitive Analysis (Apr 1 - Apr 30)",
        "Conduct 25 customer interviews across enterprise segment",
        "Benchmark against Salesforce, HubSpot, and Monday.com",
        "Deliverable: 40-page market landscape report",
        "Budget allocation: $85,000",
        "Lead: Marcus Johnson, Head of Market Intelligence",
    ])

    # --- Slide 3: Phase 2 Plan (KEY SLIDE - will be duplicated) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide3, "Phase 2 Plan", [
        "Product Development Sprint (May 1 - Jun 30)",
        "Build core API integration layer for 12 partner platforms",
        "Implement real-time analytics dashboard with D3.js",
        "QA testing: 3 rounds of regression + load testing",
        "Budget allocation: $240,000",
        "Lead: Priya Sharma, Director of Engineering",
        "Milestone: Beta release to 50 pilot customers by Jun 15",
    ])

    # --- Slide 4: Phase 3 Plan ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide4, "Phase 3 Plan", [
        "Go-to-Market Launch (Jul 1 - Aug 15)",
        "Execute multi-channel marketing campaign across 6 platforms",
        "Host 3 webinars targeting mid-market decision makers",
        "Sales enablement: update pitch deck and objection handling",
        "Budget allocation: $175,000",
        "Lead: Derek Williams, CMO",
    ])

    # --- Slide 5: Budget Overview ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide5, "Budget Overview", [
        "Total Project Budget: $620,000",
        "Phase 1 - Research: $85,000 (13.7%)",
        "Phase 2 - Development: $240,000 (38.7%)",
        "Phase 3 - Launch: $175,000 (28.2%)",
        "Contingency Reserve: $120,000 (19.4%)",
        "Approved by CFO Lisa Park on March 28, 2025",
    ])

    # --- Slide 6: Team Assignments ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide6, "Team Assignments", [
        "Project Sponsor: David Kim, COO",
        "Project Manager: Sarah Chen, VP of Operations",
        "Tech Lead: Priya Sharma, Director of Engineering",
        "Marketing Lead: Derek Williams, CMO",
        "Analytics: James Rodriguez, Senior Data Scientist",
        "Design: Aisha Patel, UX Design Manager",
        "QA: Tomoko Nakamura, QA Engineering Lead",
    ])

    # --- Slide 7: Legacy Timeline (OUTDATED - to be deleted) ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide7, "Legacy Timeline", [
        "** THIS TIMELINE IS FROM THE ORIGINAL Q1 PROPOSAL **",
        "Jan 15 - Feb 28: Initial scoping (COMPLETED)",
        "Mar 1 - Mar 31: Vendor selection (SUPERSEDED)",
        "Apr 1 - May 15: Development v1 (REVISED)",
        "May 16 - Jun 30: Testing (REVISED)",
        "Note: This slide reflects the original plan before restructuring",
    ])

    # --- Slide 8: Next Steps ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    add_title_and_body(slide8, "Next Steps", [
        "Immediate Actions (This Week):",
        "1. Finalize vendor contracts for cloud infrastructure",
        "2. Schedule kickoff meeting with all phase leads",
        "3. Set up project tracking in Jira (Project: HORIZON)",
        "4. Distribute updated RACI matrix to stakeholders",
        "Review cadence: Bi-weekly steering committee updates",
        "Next milestone review: April 14, 2025",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
