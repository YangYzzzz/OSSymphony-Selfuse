"""
Reward Script: Duplicate slide 3, move to position 5, rename title, delete slide 7
Task ID: impress_gf5_002
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.30): Slide 5 title is 'Phase 2 Summary'
  - Component 2 (0.25): Slide 5 body content duplicates slide 3 body content
  - Component 3 (0.15): Slide 3 title 'Phase 2 Plan' intact AND slide 5 is 'Phase 2 Summary' (compound)
  - Component 4 (0.30): 'Legacy Timeline' slide is removed from entire presentation
"""

import os
from pptx import Presentation

WORKDIR = '/home/user'
TASK_ID = 'impress_gf5_002'


def get_slide_title(slide):
    """Extract the first non-empty text from a slide (treated as title)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    return t
    return None


def get_slide_body_texts(slide):
    """Extract all non-empty text lines from a slide, excluding the first (title)."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    # Skip the first text (title), return body texts
    return texts[1:] if len(texts) > 1 else []


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    slides = list(prs.slides)
    num_slides = len(slides)
    print(f"INFO: Presentation has {num_slides} slides")

    # Component 1: Slide 5 has title 'Phase 2 Summary' (0.30 points)
    # Initial slide 5 is 'Budget Overview', so this FAILS on initial
    try:
        if num_slides >= 5:
            slide5_title = get_slide_title(slides[4])  # 0-indexed
            if slide5_title and slide5_title.strip() == 'Phase 2 Summary':
                print(f"PASS: Component 1 — Slide 5 title is 'Phase 2 Summary' (0.30 pts)")
                total_score += 0.30
            else:
                print(f"FAIL: Component 1 — Slide 5 title expected 'Phase 2 Summary', found: '{slide5_title}'")
        else:
            print(f"FAIL: Component 1 — Presentation has fewer than 5 slides ({num_slides})")
    except Exception as e:
        print(f"ERROR: Component 1 — {e}")

    # Component 2: Slide 5 body content duplicates Slide 3 body content (0.25 points)
    # Initial slide 5 is 'Budget Overview' with different body, so this FAILS on initial
    try:
        if num_slides >= 5:
            slide3_body = get_slide_body_texts(slides[2])  # 0-indexed
            slide5_body = get_slide_body_texts(slides[4])
            if slide3_body and slide5_body and slide3_body == slide5_body:
                print(f"PASS: Component 2 — Slide 5 body matches Slide 3 body ({len(slide3_body)} items) (0.25 pts)")
                total_score += 0.25
            else:
                print(f"FAIL: Component 2 — Slide 5 body does not match Slide 3 body")
                print(f"  Slide 3 body ({len(slide3_body)} items): {slide3_body[:3]}...")
                print(f"  Slide 5 body ({len(slide5_body)} items): {slide5_body[:3]}...")
        else:
            print(f"FAIL: Component 2 — Presentation has fewer than 5 slides")
    except Exception as e:
        print(f"ERROR: Component 2 — {e}")

    # Component 3: Slide 3 title is 'Phase 2 Plan' AND Slide 5 is 'Phase 2 Summary' (0.15 points)
    # The compound check ensures the original slide 3 is intact while the duplicate exists
    # On initial, slide 5 is 'Budget Overview', so this FAILS on initial
    try:
        if num_slides >= 5:
            slide3_title = get_slide_title(slides[2])
            slide5_title = get_slide_title(slides[4])
            if (slide3_title and slide3_title.strip() == 'Phase 2 Plan' and
                    slide5_title and slide5_title.strip() == 'Phase 2 Summary'):
                print(f"PASS: Component 3 — Slide 3 is 'Phase 2 Plan' AND Slide 5 is 'Phase 2 Summary' (0.15 pts)")
                total_score += 0.15
            else:
                print(f"FAIL: Component 3 — Slide 3 title: '{slide3_title}', Slide 5 title: '{slide5_title}'")
        else:
            print(f"FAIL: Component 3 — Presentation has fewer than 5 slides")
    except Exception as e:
        print(f"ERROR: Component 3 — {e}")

    # Component 4: 'Legacy Timeline' slide is removed (0.30 points)
    # Initial has 'Legacy Timeline' at slide 7, so this FAILS on initial
    try:
        all_titles = []
        for i, slide in enumerate(slides):
            title = get_slide_title(slide)
            all_titles.append(title)

        legacy_found = any(t and 'Legacy Timeline' in t for t in all_titles)
        if not legacy_found:
            print(f"PASS: Component 4 — 'Legacy Timeline' slide has been deleted (0.30 pts)")
            total_score += 0.30
        else:
            legacy_positions = [i+1 for i, t in enumerate(all_titles) if t and 'Legacy Timeline' in t]
            print(f"FAIL: Component 4 — 'Legacy Timeline' still present at slide(s): {legacy_positions}")
    except Exception as e:
        print(f"ERROR: Component 4 — {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
file_path = f'{WORKDIR}/project_slides.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
