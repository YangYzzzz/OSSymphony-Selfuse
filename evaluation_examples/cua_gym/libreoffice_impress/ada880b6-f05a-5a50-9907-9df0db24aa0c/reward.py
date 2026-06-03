"""
FINAL REWARD SCRIPT - SUCCESS
Task: Hey, quick LibreOffice Impress tweak: I'd like the slide I’m currently editing to have a solid #800080 background, and could you also take the exact text that’s in the Title placeholder and drop it into that slide’s Notes pane for me?
Generated: 2025-09-10 15:20:19
Status: success
Model: azure-o3
Total Steps: 3
"""

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
import os

def _rgb_from_hex(hex_string: str):
    """Convert a hex color string (e.g. '800080') to an (R, G, B) tuple."""
    hex_string = hex_string.strip().lstrip('#')
    if len(hex_string) != 6:
        raise ValueError("Hex color must be 6 characters long.")
    r = int(hex_string[0:2], 16)
    g = int(hex_string[2:4], 16)
    b = int(hex_string[4:6], 16)
    return r, g, b

def _slide_has_solid_bg_color(slide, target_hex: str) -> bool:
    """Return True if the slide background is a SOLID fill of target_hex."""
    try:
        target_rgb = _rgb_from_hex(target_hex)
        fill = slide.background.fill
        if fill.type != MSO_FILL.SOLID:
            return False  # Not a solid fill
        rgb = fill.fore_color.rgb  # pptx.dml.color.RGBColor -> behaves like a tuple
        if rgb is None:
            return False
        return (rgb[0], rgb[1], rgb[2]) == target_rgb
    except Exception:
        # Any problem means verification fails for this slide
        return False

def _get_title_text(slide):
    """Return text from the slide's title placeholder (stripped)."""
    try:
        title_shape = slide.shapes.title
    except Exception:
        return None
    if title_shape is None or not hasattr(title_shape, "text"):
        return None
    text = title_shape.text.strip()
    return text if text else None

def _get_notes_text(slide):
    """Return concatenated text (\n-separated) from all shapes on the notes pane."""
    try:
        notes_slide = slide.notes_slide  # May raise if notes part missing
    except Exception:
        return ""
    if notes_slide is None:
        return ""
    texts = []
    for shp in notes_slide.shapes:
        if hasattr(shp, "text") and shp.text:
            texts.append(shp.text.strip())
    return "\n".join(texts)

def verify_impress_tweak_task(file_path: str) -> float:
    """Verify that:
    1. Some slide has a SOLID background colour #800080, **and**
    2. Exactly that slide's Notes pane contains the *exact* text from its Title placeholder.

    Scoring (progressive):
        - 0.5 for at least one slide with correct background colour
        - 0.5 for at least one slide whose notes contain its title text
        - 1.0 only if the SAME slide satisfies BOTH requirements
    """
    print(f"Verifying presentation: {file_path}")

    # Basic file validation (no points awarded for existence!)
    if not os.path.exists(file_path):
        print("✗ File not found.")
        return 0.0

    # Attempt to load the presentation
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        print(f"✗ Failed to load presentation: {exc}")
        return 0.0

    # Flags for progressive scoring
    bg_correct = False       # Any slide has correct purple background
    notes_correct = False    # Any slide's notes contain its title text
    full_success = False     # One slide satisfies both simultaneously

    # Inspect each slide
    for idx, slide in enumerate(prs.slides, start=1):
        slide_ok_bg = _slide_has_solid_bg_color(slide, "800080")
        slide_title = _get_title_text(slide)
        slide_notes = _get_notes_text(slide)
        slide_ok_notes = bool(slide_title and slide_title in slide_notes)

        # Debug prints for transparency
        if slide_ok_bg:
            print(f"✓ Slide {idx}: solid background colour #800080 confirmed.")
        if slide_ok_notes:
            print(f"✓ Slide {idx}: notes pane contains title text → '{slide_title}'.")
        if slide_ok_bg and slide_ok_notes:
            print(f"✓ Slide {idx}: BOTH background and notes requirements met.")
            full_success = True

        # Aggregate progressive flags
        bg_correct = bg_correct or slide_ok_bg
        notes_correct = notes_correct or slide_ok_notes

    # Compute progressive score
    if full_success:
        final_score = 1.0
    else:
        final_score = 0.0
        if bg_correct:
            final_score += 0.5
        if notes_correct:
            final_score += 0.5

    # Ensure score never exceeds 1.0
    final_score = min(final_score, 1.0)

    print(f"REWARD: {final_score}")
    return final_score


# ----------------- Script Entry Point for Auto-Grading -----------------
if __name__ == "__main__":
    pptx_path = "/home/user/hey_quick_libreoffice_impress_tweak_id_like_the_slide_im_currently_editing_to_have_a_solid_800080_ba_golden.pptx"
    verify_impress_tweak_task(pptx_path)

