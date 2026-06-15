"""
Reward Script: Add presenter notes to slides 2, 3, and 4
Task ID: impress_tm_057
Domain: libreoffice_impress
Scoring:
  Component 1: Slide 2 notes match expected text (0.35 pts)
  Component 2: Slide 3 notes match expected text (0.35 pts)
  Component 3: Slide 4 notes match expected text (0.30 pts)
"""

import os
import time

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_057'

# Expected notes for each slide (1-indexed)
EXPECTED_NOTES = {
    2: 'Explain the historical context, 5 minutes',
    3: 'Show the primary source documents',
    4: 'Open floor for questions, 10 minutes',
}

# Scoring weights per slide
WEIGHTS = {
    2: 0.35,
    3: 0.35,
    4: 0.30,
}


def persist_app_state(domain: str):
    """Attempt to save any unsaved LibreOffice edits via Ctrl+S."""
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def get_slide_notes(slide):
    """Safely retrieve notes text from a slide without side effects on disk."""
    try:
        # Check if notes slide exists in the XML before accessing
        # (accessing .notes_slide creates one in memory, but we don't save)
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        return notes_text
    except Exception:
        return ""


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 4 slides
    if len(prs.slides) < 4:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 4")
        print("REWARD: 0.0")
        return 0.0

    # Component 1: Slide 2 notes match expected text (0.35 pts)
    try:
        slide2_notes = get_slide_notes(prs.slides[1])  # 0-indexed
        expected2 = EXPECTED_NOTES[2]
        if slide2_notes == expected2:
            print(f"PASS: Component 1 - Slide 2 notes exact match (0.35 pts)")
            total_score += 0.35
        elif expected2.lower() in slide2_notes.lower() and len(slide2_notes) > 0:
            print(f"PARTIAL: Component 1 - Slide 2 notes contain expected but not exact (0.2 pts). Found: {repr(slide2_notes)}")
            total_score += 0.20  # partial credit for containing expected text
        else:
            print(f"FAIL: Component 1 - Slide 2 notes mismatch")
            print(f"  Expected: {repr(expected2)}")
            print(f"  Found:    {repr(slide2_notes)}")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: Slide 3 notes match expected text (0.35 pts)
    try:
        slide3_notes = get_slide_notes(prs.slides[2])  # 0-indexed
        expected3 = EXPECTED_NOTES[3]
        if slide3_notes == expected3:
            print(f"PASS: Component 2 - Slide 3 notes exact match (0.35 pts)")
            total_score += 0.35
        elif expected3.lower() in slide3_notes.lower() and len(slide3_notes) > 0:
            print(f"PARTIAL: Component 2 - Slide 3 notes contain expected but not exact (0.2 pts). Found: {repr(slide3_notes)}")
            total_score += 0.20  # partial credit for containing expected text
        else:
            print(f"FAIL: Component 2 - Slide 3 notes mismatch")
            print(f"  Expected: {repr(expected3)}")
            print(f"  Found:    {repr(slide3_notes)}")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: Slide 4 notes match expected text (0.30 pts)
    try:
        slide4_notes = get_slide_notes(prs.slides[3])  # 0-indexed
        expected4 = EXPECTED_NOTES[4]
        if slide4_notes == expected4:
            print(f"PASS: Component 3 - Slide 4 notes exact match (0.30 pts)")
            total_score += 0.30
        elif expected4.lower() in slide4_notes.lower() and len(slide4_notes) > 0:
            print(f"PARTIAL: Component 3 - Slide 4 notes contain expected but not exact (0.15 pts). Found: {repr(slide4_notes)}")
            total_score += 0.15  # partial credit for containing expected text
        else:
            print(f"FAIL: Component 3 - Slide 4 notes mismatch")
            print(f"  Expected: {repr(expected4)}")
            print(f"  Found:    {repr(slide4_notes)}")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persist any unsaved GUI state before verification
persist_app_state("libreoffice_impress")

# Run verification
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
