"""
Initial Setup: Add presenter notes to slides 2, 3, and 4 of a history lecture
Task ID: impress_tm_057
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_tm_057'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_text_to_placeholder(slide, idx, text, font_size=Pt(18), bold=False):
    """Add text to a placeholder by index if it exists."""
    if idx in [ph.placeholder_format.idx for ph in slide.placeholders]:
        ph = slide.placeholders[idx]
        ph.text = text
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = font_size
            run.font.bold = bold


def add_bullet_content(slide, idx, items, font_size=Pt(16)):
    """Add bulleted content to a placeholder."""
    if idx not in [ph.placeholder_format.idx for ph in slide.placeholders]:
        return
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.level = 0
        for run in p.runs:
            run.font.size = font_size


def create_initial():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "World History: The Age of Exploration"
    slide1.placeholders[1].text = "Professor Elena Rodriguez\nDepartment of History\nSpring 2025"

    # --- Slide 2: Historical Context (NO notes) ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Historical Context"
    add_bullet_content(slide2, 1, [
        "European maritime expansion began in the 15th century",
        "Portugal and Spain led early exploration efforts",
        "The fall of Constantinople in 1453 disrupted trade routes",
        "Advances in navigation technology enabled longer voyages",
        "Economic motivations: spices, gold, and new trade routes",
    ])

    # --- Slide 3: Primary Source Documents (NO notes) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Primary Source Documents"
    add_bullet_content(slide3, 1, [
        "Columbus's diary entries (1492-1493)",
        "Treaty of Tordesillas (1494) - Division of the New World",
        "Letters from Vasco da Gama to King Manuel I",
        "Aztec codices depicting first contact",
        "Magellan's expedition logs (1519-1522)",
    ])

    # --- Slide 4: Discussion & Questions (NO notes) ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Discussion & Questions"
    add_bullet_content(slide4, 1, [
        "How did the Age of Exploration reshape global trade?",
        "What were the consequences for indigenous populations?",
        "Compare Portuguese and Spanish approaches to colonization",
        "How do primary sources challenge traditional narratives?",
        "What lessons can we draw for understanding modern globalization?",
    ])

    # --- Slide 5: Key Takeaways ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Key Takeaways"
    add_bullet_content(slide5, 1, [
        "Multiple factors drove European expansion beyond simple curiosity",
        "Primary sources provide nuanced perspectives often missing from textbooks",
        "The impact of exploration was transformative for all civilizations involved",
        "Understanding historical context is essential for interpreting events",
    ])

    # --- Slide 6: References ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "References"
    add_bullet_content(slide6, 1, [
        "Fernandez-Armesto, F. (2006). Pathfinders: A Global History of Exploration",
        "Restall, M. (2003). Seven Myths of the Spanish Conquest",
        "Subrahmanyam, S. (1997). The Career and Legend of Vasco da Gama",
        "Crosby, A. (2003). The Columbian Exchange: Biological and Cultural Consequences",
        "Schwartz, S. (2000). Victors and Vanquished: Spanish and Nahua Views of the Conquest",
    ], font_size=Pt(14))

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
