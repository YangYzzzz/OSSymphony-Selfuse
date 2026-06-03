"""
Reward Script: Convert bulleted text on slide 3 into chevron arrow SmartArt-style diagram
Task ID: impress_teach_068
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30) - Five chevron shapes exist on slide 3 with correct labels
  Component 2 (0.25) - Chevrons arranged horizontally (same top, increasing left)
  Component 3 (0.25) - First chevron is light blue (#BBDEFB), last is dark blue (#1565C0)
  Component 4 (0.20) - Intermediate chevrons have colors progressing from light to dark
"""

import os

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_068'

EXPECTED_LABELS = ['Plan', 'Design', 'Develop', 'Test', 'Deploy']
FIRST_COLOR = 'BBDEFB'   # light blue
LAST_COLOR = '1565C0'    # dark blue


def hex_to_rgb(h):
    """Convert hex string to (R, G, B) tuple."""
    h = h.lstrip('#')
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def color_distance(c1, c2):
    """Euclidean distance between two RGB tuples."""
    return sum((a - b) ** 2 for a, b in zip(c1, c2)) ** 0.5


def get_shape_fill_color(shape):
    """Extract solid fill srgbClr hex from shape XML. Returns uppercase hex string or None."""
    try:
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
        # Try p:spPr first, then a:spPr
        spPr = shape._element.find(f'{ns_p}spPr')
        if spPr is None:
            spPr = shape._element.find(f'{ns_a}spPr')
        if spPr is None:
            spPr = shape._element.find(f'.//{ns_a}spPr')
        if spPr is None:
            return None
        solidFill = spPr.find(f'{ns_a}solidFill')
        if solidFill is not None:
            srgb = solidFill.find(f'{ns_a}srgbClr')
            if srgb is not None:
                return srgb.get('val', '').upper()
        # Also check gradient fill — get the first stop for approximation
        gradFill = spPr.find(f'{ns_a}gradFill')
        if gradFill is not None:
            gsLst = gradFill.find(f'{ns_a}gsLst')
            if gsLst is not None:
                stops = gsLst.findall(f'{ns_a}gs')
                if stops:
                    srgb = stops[0].find(f'{ns_a}srgbClr')
                    if srgb is not None:
                        return srgb.get('val', '').upper()
    except Exception:
        pass
    return None


def get_shape_geometry(shape):
    """Get preset geometry type from shape XML. Returns string like 'chevron' or None."""
    try:
        ns_a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        prstGeom = shape._element.find(f'.//{ns_a}prstGeom')
        if prstGeom is not None:
            return prstGeom.get('prst')
    except Exception:
        pass
    return None


def persist_app_state(domain):
    """Save any unsaved GUI state before verification."""
    import time
    os.environ["DISPLAY"] = ":0"
    if domain in {"libreoffice_calc", "libreoffice_writer", "libreoffice_impress"}:
        try:
            import pyautogui
            pyautogui.hotkey("ctrl", "s")
            time.sleep(0.8)
            print(f"PERSIST: ctrl+s sent for {domain}")
        except Exception as e:
            print(f"PERSIST_WARN: save hook failed: {e}")


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: presentation must have at least 3 slides
    if len(prs.slides) < 3:
        print(f"FAIL: Presentation has {len(prs.slides)} slides, need at least 3")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[2]  # slide 3 (0-indexed)

    # Collect chevron shapes on slide 3
    chevron_shapes = []
    for shape in slide.shapes:
        geom = get_shape_geometry(shape)
        # Accept chevron or homePlate (similar arrow shapes)
        if geom in ('chevron', 'homePlate') and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                chevron_shapes.append(shape)

    # Also check for any AUTO_SHAPE with arrow-like geometry containing expected labels
    if len(chevron_shapes) < 5:
        for shape in slide.shapes:
            if str(shape.shape_type) == 'AUTO_SHAPE (1)' and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text in EXPECTED_LABELS:
                    already = any(s.text_frame.text.strip() == text for s in chevron_shapes)
                    if not already:
                        chevron_shapes.append(shape)

    print(f"INFO: Found {len(chevron_shapes)} candidate chevron/arrow shapes on slide 3")

    # Component 1: Five chevron shapes with correct labels (0.30 points)
    try:
        found_labels = [s.text_frame.text.strip() for s in chevron_shapes]
        correct_count = sum(1 for label in EXPECTED_LABELS if label in found_labels)

        if correct_count == 5 and len(chevron_shapes) == 5:
            print(f"PASS: Component 1 -- All 5 chevron shapes found with correct labels: {found_labels} (0.30 pts)")
            total_score += 0.30
        elif correct_count >= 3:
            partial = 0.30 * (correct_count / 5)
            print(f"PARTIAL: Component 1 -- {correct_count}/5 labels found: {found_labels} ({partial:.2f} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 1 -- Only {correct_count}/5 expected labels found. Shapes: {found_labels}")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Sort chevrons by horizontal position for subsequent checks
    chevron_shapes.sort(key=lambda s: s.left)

    # Component 2: Chevrons arranged horizontally (0.25 points)
    # All at same vertical position (top), with increasing left positions
    try:
        if len(chevron_shapes) >= 5:
            tops = [s.top for s in chevron_shapes]
            lefts = [s.left for s in chevron_shapes]
            # Check same top (within tolerance)
            top_spread = max(tops) - min(tops)
            same_top = top_spread < 100000  # ~0.1 inch tolerance

            # Check increasing left positions
            increasing_left = all(lefts[i] < lefts[i+1] for i in range(len(lefts)-1))

            if same_top and increasing_left:
                print(f"PASS: Component 2 -- Chevrons horizontally arranged (top spread={top_spread} EMU, lefts increasing) (0.25 pts)")
                total_score += 0.25
            elif increasing_left:
                print(f"PARTIAL: Component 2 -- Left positions increase but tops vary (spread={top_spread} EMU) (0.12 pts)")
                total_score += 0.12
            else:
                print(f"FAIL: Component 2 -- Not horizontally arranged. Tops={tops}, Lefts={lefts}")
        else:
            print(f"FAIL: Component 2 -- Not enough chevron shapes ({len(chevron_shapes)}) for arrangement check")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: First chevron is light blue (#BBDEFB), last is dark blue (#1565C0) (0.25 points)
    try:
        if len(chevron_shapes) >= 5:
            first_color = get_shape_fill_color(chevron_shapes[0])
            last_color = get_shape_fill_color(chevron_shapes[4])
            print(f"INFO: First chevron color: {first_color}, Last chevron color: {last_color}")

            first_ok = False
            last_ok = False

            if first_color is not None:
                dist_first = color_distance(hex_to_rgb(first_color), hex_to_rgb(FIRST_COLOR))
                first_ok = dist_first < 30  # allow small tolerance
                if first_ok:
                    print(f"PASS: First chevron color #{first_color} matches expected #{FIRST_COLOR} (dist={dist_first:.1f})")
                else:
                    print(f"FAIL: First chevron color #{first_color} too far from #{FIRST_COLOR} (dist={dist_first:.1f})")

            if last_color is not None:
                dist_last = color_distance(hex_to_rgb(last_color), hex_to_rgb(LAST_COLOR))
                last_ok = dist_last < 30
                if last_ok:
                    print(f"PASS: Last chevron color #{last_color} matches expected #{LAST_COLOR} (dist={dist_last:.1f})")
                else:
                    print(f"FAIL: Last chevron color #{last_color} too far from #{LAST_COLOR} (dist={dist_last:.1f})")

            if first_ok and last_ok:
                print(f"PASS: Component 3 -- Endpoint colors correct (0.25 pts)")
                total_score += 0.25
            elif first_ok or last_ok:
                print(f"PARTIAL: Component 3 -- One endpoint color correct (0.12 pts)")
                total_score += 0.12
            else:
                print(f"FAIL: Component 3 -- Neither endpoint color matches")
        else:
            print(f"FAIL: Component 3 -- Not enough chevron shapes for color check")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: Intermediate colors progress from light to dark (0.20 points)
    # The blue channel should decrease (or stay constant) from first to last
    try:
        if len(chevron_shapes) >= 5:
            colors = []
            for s in chevron_shapes:
                c = get_shape_fill_color(s)
                if c is not None:
                    colors.append(hex_to_rgb(c))
                else:
                    colors.append(None)

            if all(c is not None for c in colors):
                # Check that colors progressively get darker (lower total RGB or lower blue)
                # For a blue gradient, the overall brightness should decrease
                brightnesses = [sum(c) for c in colors]
                monotonic_decreasing = all(
                    brightnesses[i] >= brightnesses[i+1] - 20  # small tolerance
                    for i in range(len(brightnesses)-1)
                )

                # Also check that first is lighter than last
                first_brighter = brightnesses[0] > brightnesses[-1]

                if monotonic_decreasing and first_brighter:
                    print(f"PASS: Component 4 -- Colors progress from light to dark. Brightnesses: {brightnesses} (0.20 pts)")
                    total_score += 0.20
                elif first_brighter:
                    print(f"PARTIAL: Component 4 -- First lighter than last but not strictly monotonic. Brightnesses: {brightnesses} (0.10 pts)")
                    total_score += 0.10
                else:
                    print(f"FAIL: Component 4 -- Color progression incorrect. Brightnesses: {brightnesses}")
            else:
                print(f"FAIL: Component 4 -- Could not extract all colors: {colors}")
        else:
            print(f"FAIL: Component 4 -- Not enough chevron shapes for gradient check")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Entry point
persist_app_state("libreoffice_impress")

file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
