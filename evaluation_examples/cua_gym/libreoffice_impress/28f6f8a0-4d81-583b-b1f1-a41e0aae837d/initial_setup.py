"""
Initial Setup: Software Engineering presentation with bulleted list on slide 3
Task ID: impress_teach_068
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_teach_068'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Software Engineering Fundamentals"
    slide1.placeholders[1].text = "A Comprehensive Guide to Modern Software Development"

    # --- Slide 2: Methodologies Overview ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Development Methodologies"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Agile frameworks emphasize iterative development and continuous feedback"
    p2a = tf2.add_paragraph()
    p2a.text = "Waterfall follows a sequential design process through defined phases"
    p2a.level = 0
    p2b = tf2.add_paragraph()
    p2b.text = "DevOps bridges development and operations for faster delivery"
    p2b.level = 0
    p2c = tf2.add_paragraph()
    p2c.text = "Lean development minimizes waste and maximizes value"
    p2c.level = 0

    # --- Slide 3: SDLC Phases (the bulleted list to be converted) ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Software Development Lifecycle"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Plan"
    for item in ["Design", "Develop", "Test", "Deploy"]:
        p = tf3.add_paragraph()
        p.text = item
        p.level = 0

    # --- Slide 4: Requirements Engineering ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Requirements Engineering"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Stakeholder interviews and workshops gather business needs"
    p4a = tf4.add_paragraph()
    p4a.text = "User stories define features from the end-user perspective"
    p4a.level = 0
    p4b = tf4.add_paragraph()
    p4b.text = "Requirements traceability ensures complete coverage"
    p4b.level = 0

    # --- Slide 5: Architecture Patterns ---
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "Architecture Patterns"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Microservices decompose applications into loosely coupled services"
    p5a = tf5.add_paragraph()
    p5a.text = "Event-driven architecture enables reactive, scalable systems"
    p5a.level = 0
    p5b = tf5.add_paragraph()
    p5b.text = "Monolithic architecture suits smaller applications with simple deployment"
    p5b.level = 0

    # --- Slide 6: Testing Strategies ---
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "Testing Strategies"
    tf6 = slide6.placeholders[1].text_frame
    tf6.text = "Unit tests verify individual components in isolation"
    p6a = tf6.add_paragraph()
    p6a.text = "Integration tests validate interactions between modules"
    p6a.level = 0
    p6b = tf6.add_paragraph()
    p6b.text = "End-to-end tests confirm complete user workflows"
    p6b.level = 0
    p6c = tf6.add_paragraph()
    p6c.text = "Performance tests ensure the system meets SLA requirements"
    p6c.level = 0

    # --- Slide 7: CI/CD Pipeline ---
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "CI/CD Pipeline"
    tf7 = slide7.placeholders[1].text_frame
    tf7.text = "Continuous integration merges code changes multiple times a day"
    p7a = tf7.add_paragraph()
    p7a.text = "Automated builds trigger on each commit to the repository"
    p7a.level = 0
    p7b = tf7.add_paragraph()
    p7b.text = "Continuous deployment pushes validated changes to production"
    p7b.level = 0

    # --- Slide 8: Key Metrics ---
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "Key Performance Metrics"
    tf8 = slide8.placeholders[1].text_frame
    tf8.text = "Lead time: average 4.2 days from commit to production"
    p8a = tf8.add_paragraph()
    p8a.text = "Deployment frequency: 12 releases per sprint"
    p8a.level = 0
    p8b = tf8.add_paragraph()
    p8b.text = "Change failure rate: target below 5%"
    p8b.level = 0
    p8c = tf8.add_paragraph()
    p8c.text = "Mean time to recovery: under 30 minutes"
    p8c.level = 0

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
