"""
Reward Script: Customer Journey Map on Slide 7
Task ID: impress_sales_067
Domain: libreoffice_impress
Scoring:
  Component 1 (0.30): 5 rounded rectangles with stage names
  Component 2 (0.25): 5 touchpoint text boxes below stages (2 touchpoints each)
  Component 3 (0.25): 5 sentiment indicator circles (ovals) above stages
  Component 4 (0.20): 4 arrow connectors between stages
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_067'

# Expected stage names in order
STAGE_NAMES = ['Awareness', 'Consideration', 'Decision', 'Onboarding', 'Advocacy']


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: Must have at least 7 slides
    if len(prs.slides) < 7:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 7")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[6]  # Slide 7 (0-indexed)

    # Collect shapes by type for analysis
    rounded_rects = []
    ovals = []
    text_boxes = []
    connectors = []

    for shape in slide.shapes:
        # Skip the original title placeholder and "Customer Journey" text box
        # These exist in both initial and golden
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            try:
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
                    text = ""
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                    rounded_rects.append({'shape': shape, 'text': text})
                elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                    text = ""
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                    ovals.append({'shape': shape, 'text': text})
            except Exception:
                pass

        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text = ""
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
            # Only consider non-title text boxes (skip "Customer Journey" title)
            if text.lower() != "customer journey":
                text_boxes.append({'shape': shape, 'text': text})

        elif shape.shape_type in (MSO_SHAPE_TYPE.LINE,):
            connectors.append(shape)

    # Also check for freeform/connector types
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        # Some connectors may be shape_type 9 (LINE) which we already handle
        # Also check for type 13 (FREEFORM) used as connectors sometimes
        pass

    print(f"Found: {len(rounded_rects)} rounded rects, {len(ovals)} ovals, "
          f"{len(text_boxes)} text boxes, {len(connectors)} connectors")

    # Component 1: 5 Rounded Rectangles with correct stage names (0.30 points)
    # Each correctly named stage earns 0.06 points
    try:
        found_stages = set()
        for rr in rounded_rects:
            for stage_name in STAGE_NAMES:
                if stage_name.lower() in rr['text'].lower():
                    found_stages.add(stage_name)

        comp1_score = len(found_stages) * 0.06
        if len(found_stages) == 5:
            print(f"PASS: Component 1 - All 5 stage rectangles found: {found_stages} (0.30 pts)")
            total_score += 0.30
        elif len(found_stages) > 0:
            print(f"PARTIAL: Component 1 - {len(found_stages)}/5 stages found: {found_stages} ({comp1_score:.2f} pts)")
            total_score += comp1_score
        else:
            print(f"FAIL: Component 1 - No stage rectangles found (expected 5 rounded rectangles with stage names)")
    except Exception as e:
        print(f"ERROR: Component 1 - {e}")

    # Component 2: 5 Touchpoint text boxes below stages (0.25 points)
    # Each text box with at least 2 touchpoints earns 0.05 points
    try:
        valid_touchpoint_boxes = 0
        for tb in text_boxes:
            text = tb['text']
            # Count lines/bullets - touchpoints are typically on separate lines
            lines = [l.strip() for l in text.split('\n') if l.strip()]
            # Also check paragraphs
            para_count = 0
            if tb['shape'].has_text_frame:
                para_count = len([p for p in tb['shape'].text_frame.paragraphs if p.text.strip()])
            num_items = max(len(lines), para_count)
            if num_items >= 2:
                valid_touchpoint_boxes += 1

        comp2_score = min(valid_touchpoint_boxes, 5) * 0.05
        if valid_touchpoint_boxes >= 5:
            print(f"PASS: Component 2 - {valid_touchpoint_boxes} touchpoint boxes with 2+ items (0.25 pts)")
            total_score += 0.25
        elif valid_touchpoint_boxes > 0:
            print(f"PARTIAL: Component 2 - {valid_touchpoint_boxes}/5 touchpoint boxes ({comp2_score:.2f} pts)")
            total_score += comp2_score
        else:
            print(f"FAIL: Component 2 - No valid touchpoint text boxes found")
    except Exception as e:
        print(f"ERROR: Component 2 - {e}")

    # Component 3: 5 Sentiment indicator circles/ovals above stages (0.25 points)
    # Each oval earns 0.05 points
    try:
        sentiment_count = len(ovals)
        comp3_score = min(sentiment_count, 5) * 0.05
        if sentiment_count >= 5:
            print(f"PASS: Component 3 - {sentiment_count} sentiment indicator ovals found (0.25 pts)")
            total_score += 0.25
        elif sentiment_count > 0:
            print(f"PARTIAL: Component 3 - {sentiment_count}/5 ovals ({comp3_score:.2f} pts)")
            total_score += comp3_score
        else:
            print(f"FAIL: Component 3 - No sentiment indicator ovals found")
    except Exception as e:
        print(f"ERROR: Component 3 - {e}")

    # Component 4: Arrow connectors between stages (0.20 points)
    # 4 connectors expected, each earns 0.05 points
    try:
        connector_count = len(connectors)
        comp4_score = min(connector_count, 4) * 0.05
        if connector_count >= 4:
            print(f"PASS: Component 4 - {connector_count} arrow connectors found (0.20 pts)")
            total_score += 0.20
        elif connector_count > 0:
            print(f"PARTIAL: Component 4 - {connector_count}/4 connectors ({comp4_score:.2f} pts)")
            total_score += comp4_score
        else:
            print(f"FAIL: Component 4 - No arrow connectors found")
    except Exception as e:
        print(f"ERROR: Component 4 - {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Persistence hook for LibreOffice Impress
def persist_app_state():
    import time
    os.environ["DISPLAY"] = ":0"
    try:
        import pyautogui
        pyautogui.hotkey("ctrl", "s")
        time.sleep(0.8)
        print("PERSIST: ctrl+s sent for libreoffice_impress")
    except Exception as e:
        print(f"PERSIST_WARN: save hook failed: {e}")


# Entry point
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    persist_app_state()
    verify_task(file_path)
