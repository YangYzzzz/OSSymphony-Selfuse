"""
Initial Setup: Customer Journey Map presentation with empty slide 7
Task ID: impress_sales_067
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_sales_067'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_title_content_slide(prs, title_text, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "Journey Pitch", "Customer Experience Strategy 2025\nPrepared by Sarah Chen, VP of CX")

    # Slide 2: Agenda
    add_title_content_slide(prs, "Agenda", [
        "1. Market Overview & Customer Landscape",
        "2. Current Customer Satisfaction Metrics",
        "3. Competitive Benchmarking",
        "4. Voice of Customer Insights",
        "5. Proposed Journey Improvements",
        "6. Customer Segmentation Analysis",
        "7. Customer Journey Map",
        "8. Implementation Roadmap",
        "9. Budget & Resource Allocation",
        "10. Next Steps & Timeline",
    ])

    # Slide 3: Market Overview
    add_title_content_slide(prs, "Market Overview & Customer Landscape", [
        "Total addressable market grew 18% YoY to $4.2B",
        "Customer base expanded from 12,400 to 15,800 accounts",
        "Enterprise segment represents 62% of revenue ($1.84B)",
        "Mid-market adoption accelerated by 24% in Q3-Q4",
        "Average contract value increased to $28,500",
        "Net Promoter Score improved from 42 to 56",
    ])

    # Slide 4: Current Satisfaction Metrics
    add_title_content_slide(prs, "Current Customer Satisfaction Metrics", [
        "Overall CSAT: 4.2/5.0 (+0.3 from last year)",
        "First Contact Resolution: 78% (target: 85%)",
        "Average Response Time: 2.4 hours (down from 4.1 hours)",
        "Customer Effort Score: 3.8/7.0 (lower is better)",
        "Retention Rate: 91% (up from 87%)",
        "Churn Rate: 9% (target: <7%)",
    ])

    # Slide 5: Competitive Benchmarking
    add_title_content_slide(prs, "Competitive Benchmarking", [
        "We rank #2 in customer satisfaction behind Acme Corp (4.5/5.0)",
        "Our onboarding time (14 days) lags industry best (7 days)",
        "Support quality rated 'excellent' by 68% of customers",
        "Feature parity achieved in 8 of 10 key categories",
        "Price-to-value perception: 7.2/10 (industry avg: 6.8)",
        "Brand trust index: 82/100 (top quartile)",
    ])

    # Slide 6: Voice of Customer
    add_title_content_slide(prs, "Voice of Customer Insights", [
        "\"The onboarding process could be more streamlined\" - recurring theme",
        "Top request: Self-service portal with real-time dashboards",
        "Pain point: Handoff between sales and customer success teams",
        "85% of enterprise clients want dedicated account managers",
        "Mobile experience rated 3.1/5.0 - significant gap identified",
        "Integration with existing tools cited as top purchase factor",
    ])

    # Slide 7: Customer Journey (TITLE ONLY - empty content)
    slide7 = add_title_only_slide(prs, "Customer Journey")

    # Slide 8: Implementation Roadmap
    add_title_content_slide(prs, "Implementation Roadmap", [
        "Phase 1 (Q1): Discovery & journey mapping workshops",
        "Phase 2 (Q2): Redesign onboarding flow, launch self-service portal",
        "Phase 3 (Q3): Integrate CRM touchpoint tracking, deploy feedback loops",
        "Phase 4 (Q4): Measure improvements, iterate on low-performing stages",
        "Key milestone: All journey improvements live by September 30",
        "Executive review checkpoints: Monthly with C-suite sponsors",
    ])

    # Slide 9: Budget & Resources
    add_title_content_slide(prs, "Budget & Resource Allocation", [
        "Total investment: $1.2M across 4 quarters",
        "Technology platform: $450K (CRM upgrade, analytics tools)",
        "Staffing: 3 FTEs + 2 contractors for 9 months ($520K)",
        "Training & change management: $130K",
        "Contingency reserve: $100K (8% of total budget)",
        "Expected ROI: 3.2x within 18 months based on retention uplift",
    ])

    # Slide 10: Next Steps
    add_title_content_slide(prs, "Next Steps & Timeline", [
        "Secure executive sponsorship and budget approval by Jan 15",
        "Kick off cross-functional journey mapping workshop (Jan 22-24)",
        "Complete customer interview program (30 interviews by Feb 28)",
        "Present revised journey map to leadership (March 10)",
        "Begin Phase 1 implementation (March 17)",
        "Monthly progress reports to steering committee",
    ])

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
