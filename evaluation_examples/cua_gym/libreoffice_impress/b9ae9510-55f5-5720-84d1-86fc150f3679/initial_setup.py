"""
Initial Setup: Create a 10-slide ProductLaunch_X1.pptx presentation.
Slide 8 has title 'Live Demo Schedule' but empty content area.
Task ID: impress_ps_015
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_015'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def add_title_slide(prs, title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = subtitle_text
    return slide


def add_content_slide(prs, title_text, bullet_items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, item in enumerate(bullet_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    return slide


def add_title_only_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
    return slide


def create_initial():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(prs, "Product Launch X1", "Next-Generation Smart Device\nQ3 2026 Launch Strategy")

    # Slide 2: Agenda
    add_content_slide(prs, "Agenda", [
        "Market Opportunity & Research",
        "Product Features & Specifications",
        "Competitive Analysis",
        "Pricing Strategy",
        "Go-to-Market Plan",
        "Live Demo Schedule",
        "Next Steps & Timeline",
    ])

    # Slide 3: Market Analysis
    add_content_slide(prs, "Market Analysis", [
        "Total addressable market: $4.2B (growing 18% YoY)",
        "Consumer electronics segment expanding rapidly",
        "Key demographics: 25-45 professionals, tech-savvy",
        "78% of target segment willing to pay premium for innovation",
        "Major competitors: TechCo Pro ($899), InnoDevice Max ($1,049)",
    ])

    # Slide 4: Product Features
    add_content_slide(prs, "Product Features — X1 Overview", [
        "12-core Neural Processing Unit (NPU) at 3.8 GHz",
        "8K HDR display with 120Hz refresh rate",
        "Advanced thermal management with vapor chamber cooling",
        "AI-powered voice assistant with offline mode",
        "72-hour battery life with fast charging (0-80% in 15 min)",
    ])

    # Slide 5: Competitive Advantage
    add_content_slide(prs, "Competitive Advantage", [
        "35% better performance vs TechCo Pro in benchmark tests",
        "Only device with offline AI capabilities in this price range",
        "Patented cooling system allows sustained peak performance",
        "Ecosystem integration with 200+ smart home devices",
        "3-year warranty with premium support included",
    ])

    # Slide 6: Pricing Strategy
    add_content_slide(prs, "Pricing Strategy", [
        "MSRP: $949 (Standard) / $1,199 (Pro Edition)",
        "Early adopter discount: 15% off for first 30 days",
        "Channel partner margin: 22%",
        "Projected break-even: Month 8 post-launch",
        "Lifetime value per customer: $2,400 (accessories + services)",
    ])

    # Slide 7: Go-to-Market Plan
    add_content_slide(prs, "Go-to-Market Plan", [
        "Phase 1 (Week 1-4): Influencer seeding & tech review embargoes",
        "Phase 2 (Week 5-8): Public launch event in San Francisco",
        "Phase 3 (Week 9-16): Retail rollout across 12 countries",
        "Digital campaign: $2.5M budget across social, search, display",
        "Partnership activations with major carriers & retailers",
    ])

    # Slide 8: Live Demo Schedule — TITLE ONLY, EMPTY CONTENT AREA
    add_title_only_slide(prs, "Live Demo Schedule")

    # Slide 9: Next Steps
    add_content_slide(prs, "Next Steps & Timeline", [
        "Finalize packaging design — by June 15, 2026",
        "Complete FCC/CE certification — by July 1, 2026",
        "Production ramp-up: 50,000 units initial batch",
        "Media review samples shipped — July 20, 2026",
        "Public launch event — August 5, 2026",
    ])

    # Slide 10: Thank You / Q&A
    s10 = add_title_slide(prs, "Thank You", "Questions & Discussion\ncontact@productx1.com")

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')
    print(f'Total slides: {len(prs.slides)}')

    # Open in LibreOffice Impress
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
