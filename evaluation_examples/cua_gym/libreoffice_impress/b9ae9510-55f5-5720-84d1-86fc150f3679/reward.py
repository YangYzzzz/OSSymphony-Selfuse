"""
Reward Script: Create a demo timeline slide with 4 proportional segments at position 8
Task ID: impress_ps_015
Domain: libreoffice_impress
Scoring:
  - Component 1 (0.30): 4 rectangle segments exist on slide 8 with correct label text
  - Component 2 (0.25): Segment widths are proportional to durations (2, 3, 7, 8 min)
  - Component 3 (0.20): Each segment has a different fill color
  - Component 4 (0.15): 5 time marker text boxes below the bar (0:00, 2:00, 5:00, 12:00, 20:00)
  - Component 5 (0.10): Total bar spans ~80% of slide width
"""

import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ps_015'

# Expected segment data
EXPECTED_SEGMENTS = [
    {'label': 'Unboxing', 'time': '(0:00-2:00)', 'duration': 2},
    {'label': 'Setup', 'time': '(2:00-5:00)', 'duration': 3},
    {'label': 'Performance Tests', 'time': '(5:00-12:00)', 'duration': 7},
    {'label': 'Q&A', 'time': '(12:00-20:00)', 'duration': 8},
]
TOTAL_DURATION = 20
EXPECTED_MARKERS = ['0:00', '2:00', '5:00', '12:00', '20:00']

SLIDE_INDEX = 7  # Slide 8, 0-based


def verify_task(file_path):
    """
    Verify task completion with progressive scoring.
    Returns: float between 0.0 and 1.0
    """
    total_score = 0.0

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"CRITICAL: Cannot load file {file_path}: {e}")
        print("REWARD: 0.0")
        return 0.0

    # Precondition: at least 8 slides
    if len(prs.slides) < 8:
        print(f"FAIL: Presentation has only {len(prs.slides)} slides, need at least 8")
        print("REWARD: 0.0")
        return 0.0

    slide = prs.slides[SLIDE_INDEX]
    slide_width = prs.slide_width

    # Collect rectangles (AUTO_SHAPE) and text boxes on this slide
    rectangles = []
    textboxes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rectangles.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Exclude the title text box (the pre-existing "Live Demo Schedule")
            text = shape.text.strip() if hasattr(shape, 'text') else ''
            textboxes.append(shape)

    # Sort rectangles by left position (left to right)
    rectangles.sort(key=lambda s: s.left)

    # Component 1: 4 rectangle segments with correct labels (0.30 points)
    try:
        if len(rectangles) < 4:
            print(f"FAIL: Component 1 -- Expected 4 rectangle segments, found {len(rectangles)}")
        else:
            matched = 0
            for i, seg in enumerate(EXPECTED_SEGMENTS):
                if i < len(rectangles):
                    rect_text = rectangles[i].text.replace('\x0b', '\n').replace('\n', ' ').strip()
                    # Check if both label and time range are present
                    has_label = seg['label'].lower() in rect_text.lower()
                    has_time = seg['time'] in rect_text
                    if has_label and has_time:
                        matched += 1
                        print(f"  PASS: Segment {i+1} '{seg['label']}' found with correct time range")
                    else:
                        print(f"  FAIL: Segment {i+1} expected '{seg['label']} {seg['time']}', found '{rect_text}'")
            if matched == 4:
                print(f"PASS: Component 1 -- All 4 segments have correct labels (0.30 pts)")
                total_score += 0.30
            elif matched >= 2:
                partial = 0.15
                print(f"PARTIAL: Component 1 -- {matched}/4 segments correct ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 1 -- Only {matched}/4 segments matched")
    except Exception as e:
        print(f"ERROR: Component 1 -- {e}")

    # Component 2: Segment widths proportional to durations (0.25 points)
    try:
        if len(rectangles) < 4:
            print(f"FAIL: Component 2 -- Not enough rectangles to check proportions")
        else:
            total_bar_width = sum(rectangles[i].width for i in range(4))
            if total_bar_width <= 0:
                print(f"FAIL: Component 2 -- Total bar width is 0")
            else:
                proportion_ok = 0
                for i, seg in enumerate(EXPECTED_SEGMENTS):
                    expected_ratio = seg['duration'] / TOTAL_DURATION
                    actual_ratio = rectangles[i].width / total_bar_width
                    diff = abs(actual_ratio - expected_ratio)
                    if diff <= 0.05:  # 5% tolerance
                        proportion_ok += 1
                        print(f"  PASS: Segment {i+1} ratio {actual_ratio:.3f} ~ expected {expected_ratio:.3f}")
                    else:
                        print(f"  FAIL: Segment {i+1} ratio {actual_ratio:.3f} != expected {expected_ratio:.3f} (diff={diff:.3f})")
                if proportion_ok == 4:
                    print(f"PASS: Component 2 -- All segment widths proportional (0.25 pts)")
                    total_score += 0.25
                elif proportion_ok >= 2:
                    partial = 0.12
                    print(f"PARTIAL: Component 2 -- {proportion_ok}/4 proportions correct ({partial} pts)")
                    total_score += partial
                else:
                    print(f"FAIL: Component 2 -- Only {proportion_ok}/4 proportions correct")
    except Exception as e:
        print(f"ERROR: Component 2 -- {e}")

    # Component 3: Each segment has a different fill color (0.20 points)
    try:
        if len(rectangles) < 4:
            print(f"FAIL: Component 3 -- Not enough rectangles")
        else:
            colors = []
            for i in range(4):
                shape = rectangles[i]
                fill = shape.fill
                color_str = None
                try:
                    if fill.type is not None:
                        color_str = str(fill.fore_color.rgb)
                except Exception:
                    pass
                colors.append(color_str)

            non_none_colors = [c for c in colors if c is not None]
            if len(non_none_colors) == 4 and len(set(non_none_colors)) == 4:
                print(f"PASS: Component 3 -- 4 different fill colors: {non_none_colors} (0.20 pts)")
                total_score += 0.20
            elif len(non_none_colors) >= 2 and len(set(non_none_colors)) >= 2:
                partial = 0.10
                print(f"PARTIAL: Component 3 -- {len(set(non_none_colors))} unique colors out of {len(non_none_colors)} filled ({partial} pts)")
                total_score += partial
            else:
                print(f"FAIL: Component 3 -- Colors: {colors}")
    except Exception as e:
        print(f"ERROR: Component 3 -- {e}")

    # Component 4: 5 time marker text boxes (0.15 points)
    try:
        # Time markers should be text boxes below the bar
        # Filter textboxes that contain time marker text (exclude pre-existing title)
        marker_texts_found = []
        for tb in textboxes:
            text = tb.text.strip()
            if text in EXPECTED_MARKERS:
                marker_texts_found.append(text)

        unique_markers = set(marker_texts_found)
        if len(unique_markers) == 5:
            print(f"PASS: Component 4 -- All 5 time markers found: {sorted(unique_markers)} (0.15 pts)")
            total_score += 0.15
        elif len(unique_markers) >= 3:
            partial = 0.08
            print(f"PARTIAL: Component 4 -- {len(unique_markers)}/5 time markers found ({partial} pts)")
            total_score += partial
        else:
            print(f"FAIL: Component 4 -- Only {len(unique_markers)} time markers found: {unique_markers}")
    except Exception as e:
        print(f"ERROR: Component 4 -- {e}")

    # Component 5: Total bar spans ~80% of slide width (0.10 points)
    try:
        if len(rectangles) < 4:
            print(f"FAIL: Component 5 -- Not enough rectangles")
        else:
            total_bar_width = sum(rectangles[i].width for i in range(4))
            bar_ratio = total_bar_width / slide_width
            if 0.65 <= bar_ratio <= 0.95:
                print(f"PASS: Component 5 -- Bar spans {bar_ratio:.1%} of slide width (0.10 pts)")
                total_score += 0.10
            else:
                print(f"FAIL: Component 5 -- Bar spans {bar_ratio:.1%}, expected ~80%")
    except Exception as e:
        print(f"ERROR: Component 5 -- {e}")

    final_score = min(total_score, 1.0)
    print(f"\nScore: {total_score}/1.0")
    print(f"REWARD: {final_score}")
    return final_score


# Default: test against canonical artifact path
file_path = f'{WORKDIR}/{TASK_ID}.pptx'
if not os.path.exists(file_path):
    print(f"File not found: {file_path}")
    print("REWARD: 0.0")
else:
    verify_task(file_path)
