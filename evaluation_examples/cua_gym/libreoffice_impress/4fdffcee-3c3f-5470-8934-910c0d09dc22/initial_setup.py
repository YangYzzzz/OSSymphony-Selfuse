"""
Initial Setup: Rectangle with solid fill on slide 2
Task ID: impress_ndo_049
Domain: libreoffice_impress
"""

import os
import shlex
import subprocess
import time

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

WORKDIR = '/home/user'
TASK_ID = 'impress_ndo_049'
OUTPUT = f'{WORKDIR}/{TASK_ID}.pptx'


def launch_gui(command: str, delay_sec: float = 1.0):
    """Launch GUI app on VM display without blocking script exit."""
    env = os.environ.copy()
    env["DISPLAY"] = ":0"
    subprocess.Popen(
        shlex.split(command),
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        env=env,
    )
    time.sleep(delay_sec)


def create_initial():
    prs = Presentation()
    prs.slide_width = Cm(25.4)   # standard 10 inches
    prs.slide_height = Cm(19.05) # standard 7.5 inches

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Patterns & Design Elements"
    slide1.placeholders[1].text = "Q2 2025 Visual Guidelines"

    # --- Slide 2: Contains the target rectangle ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

    # Add a title text box at the top
    txBox = slide2.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Pattern Samples"
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # The TARGET rectangle: 10cm x 6cm with solid fill #CCCCCC
    rect = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(7.7),   # centered horizontally: (25.4 - 10) / 2
        Cm(5.5),   # positioned in the middle-ish area
        Cm(10),    # width = 10cm
        Cm(6),     # height = 6cm
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    # Add a thin border
    rect.line.color.rgb = RGBColor(0x66, 0x66, 0x66)
    rect.line.width = Pt(1)

    # Add a label below the rectangle
    label = slide2.shapes.add_textbox(Cm(7.7), Cm(12), Cm(10), Cm(1.5))
    lf = label.text_frame
    lp = lf.paragraphs[0]
    lp.text = "Sample Rectangle — Solid Fill #CCCCCC"
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.runs[0]
    lr.font.name = "Arial"
    lr.font.size = Pt(14)
    lr.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # Add a small circle for visual variety
    circle = slide2.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Cm(19), Cm(6),
        Cm(3), Cm(3),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x44, 0x72, 0xC4)
    circle.line.color.rgb = RGBColor(0x2F, 0x52, 0x96)
    circle.line.width = Pt(1)

    # --- Slide 3: Additional content ---
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide3.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Color Palette Reference"
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.runs[0]
    r3.font.name = "Arial"
    r3.font.size = Pt(28)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Add a few color swatches
    colors = [
        (RGBColor(0x1A, 0x1A, 0x1A), "Charcoal #1A1A1A"),
        (RGBColor(0xF0, 0xF0, 0xF0), "Light Gray #F0F0F0"),
        (RGBColor(0xCC, 0xCC, 0xCC), "Medium Gray #CCCCCC"),
        (RGBColor(0x44, 0x72, 0xC4), "Blue #4472C4"),
    ]
    for i, (color, name) in enumerate(colors):
        swatch = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Cm(2 + i * 5.5), Cm(5),
            Cm(4), Cm(3),
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = color
        swatch.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        swatch.line.width = Pt(0.5)

        lbl = slide3.shapes.add_textbox(Cm(2 + i * 5.5), Cm(8.5), Cm(4), Cm(1))
        ltf = lbl.text_frame
        ltf.word_wrap = True
        ltp = ltf.paragraphs[0]
        ltp.text = name
        ltp.alignment = PP_ALIGN.CENTER
        ltr = ltp.runs[0]
        ltr.font.name = "Arial"
        ltr.font.size = Pt(11)
        ltr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    # --- Slide 4: Notes slide ---
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide4.shapes.add_textbox(Cm(2), Cm(2), Cm(21), Cm(14))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Design Notes"
    r4 = p4.runs[0]
    r4.font.name = "Arial"
    r4.font.size = Pt(24)
    r4.font.bold = True
    r4.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    notes_text = [
        "The pattern samples on slide 2 demonstrate fill options available for shapes.",
        "Solid fills provide a clean, professional look for business presentations.",
        "Hatching and gradient fills can add visual texture when appropriate.",
        "Always maintain sufficient contrast between fill and text colors.",
    ]
    for text in notes_text:
        np = tf4.add_paragraph()
        np.text = text
        nr = np.runs[0]
        nr.font.name = "Arial"
        nr.font.size = Pt(16)
        nr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    prs.save(OUTPUT)
    print(f'Initial file created: {OUTPUT}')

    # GUI-ready startup
    launch_gui(f'libreoffice --impress "{OUTPUT}"', delay_sec=2.0)
    print('GUI_READY: launched LibreOffice Impress with DISPLAY=:0')


create_initial()
